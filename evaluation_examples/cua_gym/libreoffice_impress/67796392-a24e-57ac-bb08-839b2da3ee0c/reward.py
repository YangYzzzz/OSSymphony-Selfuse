"""
Reward Script: Photo slideshow with transitions, music, and watermark
Task ID: impress_wf_017
Domain: libreoffice_impress
Scoring:
  Component 1: File exists with 10 slides (0.15 pts)
  Component 2: Each slide has a full-bleed image (0.20 pts)
  Component 3: Dissolve transitions with 4s auto-advance on all slides (0.25 pts)
  Component 4: Audio object on slide 1 (0.15 pts)
  Component 5: 'Sample' watermark on each slide in white with reduced opacity (0.25 pts)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_017'
FILE_PATH = os.path.join(WORKDIR, 'Desktop', 'Photo_Slideshow.pptx')

# XML namespaces used in OOXML
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'


def persist_app_state():
    """Try to save any unsaved state in LibreOffice."""
    try:
        import time
        os.environ["DISPLAY"] = ":0"
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: file must exist
    if not os.path.exists(file_path):
        print(f"CRITICAL: File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    # Load with python-pptx for shape analysis
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # ---------------------------------------------------------------
    # Component 1: File has exactly 10 slides (0.15 points)
    # This FAILS on initial (file doesn't exist) -> PASSES on golden
    # ---------------------------------------------------------------
    try:
        num_slides = len(prs.slides)
        if num_slides == 10:
            print(f"PASS: Component 1 -- 10 slides found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 -- expected 10 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # ---------------------------------------------------------------
    # Component 2: Each slide has at least one full-bleed image (0.20 points)
    # Full-bleed = image covers full slide (pos near 0,0, size near slide dimensions)
    # FAILS on initial (no file) -> PASSES on golden
    # ---------------------------------------------------------------
    try:
        slide_w = prs.slide_width
        slide_h = prs.slide_height
        slides_with_image = 0
        tolerance = 0.05  # 5% tolerance for "full bleed"

        for i, slide in enumerate(prs.slides):
            has_full_image = False
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        _ = shape.image.blob  # verify it's a real image, not audio
                    except ValueError:
                        continue  # skip audio shapes
                    # Check if image covers at least 90% of slide area
                    w_ratio = shape.width / slide_w if slide_w > 0 else 0
                    h_ratio = shape.height / slide_h if slide_h > 0 else 0
                    if w_ratio >= (1 - tolerance) and h_ratio >= (1 - tolerance):
                        has_full_image = True
                        break
            if has_full_image:
                slides_with_image += 1

        if slides_with_image == 10:
            print(f"PASS: Component 2 -- all 10 slides have full-bleed images (0.20 pts)")
            total_score += 0.20
        elif slides_with_image >= 7:
            partial = 0.20 * (slides_with_image / 10)
            print(f"PARTIAL: Component 2 -- {slides_with_image}/10 slides have full-bleed images ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 -- only {slides_with_image}/10 slides have full-bleed images")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # ---------------------------------------------------------------
    # Component 3: Dissolve transitions with 4s auto-advance on all slides (0.25 points)
    # Verified via ZIP/XML since python-pptx doesn't expose transitions
    # FAILS on initial (no file) -> PASSES on golden
    # ---------------------------------------------------------------
    try:
        slides_with_dissolve = 0
        slides_with_4s = 0

        with zipfile.ZipFile(file_path, 'r') as zf:
            for i in range(1, 11):
                fname = f'ppt/slides/slide{i}.xml'
                try:
                    with zf.open(fname) as f:
                        root = ET.parse(f).getroot()
                        tr = root.find(f'.//{{{NS_P}}}transition')
                        if tr is not None:
                            # Check for dissolve child element
                            dissolve = tr.find(f'{{{NS_P}}}dissolve')
                            if dissolve is not None:
                                slides_with_dissolve += 1

                            # Check advTm attribute for auto-advance timing
                            adv_tm = tr.get('advTm')
                            if adv_tm is not None:
                                adv_ms = int(adv_tm)
                                # Allow some tolerance: 3500-4500ms
                                if 3500 <= adv_ms <= 4500:
                                    slides_with_4s += 1
                except KeyError:
                    pass

        dissolve_score = 0.0
        if slides_with_dissolve == 10:
            dissolve_score += 0.15
            print(f"PASS: Component 3a -- all 10 slides have Dissolve transition")
        elif slides_with_dissolve >= 5:
            partial = 0.15 * (slides_with_dissolve / 10)
            dissolve_score += partial
            print(f"PARTIAL: Component 3a -- {slides_with_dissolve}/10 slides have Dissolve transition")
        else:
            print(f"FAIL: Component 3a -- only {slides_with_dissolve}/10 slides have Dissolve transition")

        if slides_with_4s == 10:
            dissolve_score += 0.10
            print(f"PASS: Component 3b -- all 10 slides have 4s auto-advance")
        elif slides_with_4s >= 5:
            partial = 0.10 * (slides_with_4s / 10)
            dissolve_score += partial
            print(f"PARTIAL: Component 3b -- {slides_with_4s}/10 slides have 4s auto-advance")
        else:
            print(f"FAIL: Component 3b -- only {slides_with_4s}/10 slides have 4s auto-advance")

        total_score += dissolve_score
        print(f"  Component 3 total: {dissolve_score:.2f} pts")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # ---------------------------------------------------------------
    # Component 4: Audio object on slide 1 (0.15 points)
    # Check slide 1 relationships for an audio reference
    # FAILS on initial (no file) -> PASSES on golden
    # ---------------------------------------------------------------
    try:
        audio_found = False
        with zipfile.ZipFile(file_path, 'r') as zf:
            rels_path = 'ppt/slides/_rels/slide1.xml.rels'
            try:
                with zf.open(rels_path) as f:
                    root = ET.parse(f).getroot()
                    for rel in root:
                        rel_type = rel.attrib.get('Type', '')
                        target = rel.attrib.get('Target', '')
                        if 'audio' in rel_type.lower() or target.lower().endswith('.mp3'):
                            audio_found = True
                            break
            except KeyError:
                pass

        if audio_found:
            print(f"PASS: Component 4 -- audio object found on slide 1 (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 -- no audio object found on slide 1")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # ---------------------------------------------------------------
    # Component 5: 'Sample' watermark text on each slide (0.25 points)
    # Check for textbox with 'Sample' text, white color, reduced opacity
    # FAILS on initial (no file) -> PASSES on golden
    # ---------------------------------------------------------------
    try:
        slides_with_watermark = 0

        with zipfile.ZipFile(file_path, 'r') as zf:
            for i in range(1, 11):
                fname = f'ppt/slides/slide{i}.xml'
                try:
                    with zf.open(fname) as f:
                        root = ET.parse(f).getroot()
                        # Find all shape trees
                        for sp in root.findall(f'.//{{{NS_P}}}sp'):
                            # Find text elements
                            texts = sp.findall(f'.//{{{NS_A}}}t')
                            text_content = ''.join(t.text or '' for t in texts).strip()
                            if 'Sample' in text_content or 'sample' in text_content.lower():
                                # Found Sample text; check for white color and opacity
                                rPr_list = sp.findall(f'.//{{{NS_A}}}rPr')
                                for rPr in rPr_list:
                                    solidFill = rPr.find(f'{{{NS_A}}}solidFill')
                                    if solidFill is not None:
                                        srgbClr = solidFill.find(f'{{{NS_A}}}srgbClr')
                                        if srgbClr is not None:
                                            color_val = srgbClr.get('val', '').upper()
                                            alpha_el = srgbClr.find(f'{{{NS_A}}}alpha')
                                            # White color
                                            if color_val == 'FFFFFF':
                                                # Check opacity (alpha val < 100000 means reduced)
                                                if alpha_el is not None:
                                                    alpha_val = int(alpha_el.get('val', '100000'))
                                                    if alpha_val < 100000:
                                                        slides_with_watermark += 1
                                                        break
                                                else:
                                                    # White text without explicit opacity still counts partially
                                                    slides_with_watermark += 1
                                                    break
                                else:
                                    continue
                                break  # Found Sample text shape, move to next slide
                except KeyError:
                    pass

        watermark_score = 0.0
        if slides_with_watermark == 10:
            watermark_score = 0.25
            print(f"PASS: Component 5 -- all 10 slides have 'Sample' watermark with white/opacity (0.25 pts)")
        elif slides_with_watermark >= 5:
            watermark_score = 0.25 * (slides_with_watermark / 10)
            print(f"PARTIAL: Component 5 -- {slides_with_watermark}/10 slides have watermark ({watermark_score:.2f} pts)")
        else:
            print(f"FAIL: Component 5 -- only {slides_with_watermark}/10 slides have 'Sample' watermark")

        total_score += watermark_score
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(FILE_PATH)
