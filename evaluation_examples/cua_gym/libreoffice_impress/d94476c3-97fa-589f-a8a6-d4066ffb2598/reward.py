"""
Reward Script: Crop image on slide 4 to top half and center it
Task ID: impress_stu_079
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.4): Bottom crop of ~50% applied to the image
  - Component 2 (0.3): Image visually becomes ~3x3 inches (top half of 3x6)
  - Component 3 (0.3): Image centered on the slide
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_079'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide4 = prs.slides[3]  # 0-indexed
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Find the picture shape on slide 4
    pic_shape = None
    for shape in slide4.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pic_shape = shape
            break

    if pic_shape is None:
        print("FAIL: No picture shape found on slide 4")
        print("REWARD: 0.0")
        return 0.0

    # Read crop values from srcRect element
    el = pic_shape._element
    src_rect = el.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srcRect')

    crop_bottom = 0
    crop_top = 0
    crop_left = 0
    crop_right = 0
    if src_rect is not None:
        crop_bottom = int(src_rect.get('b', '0'))
        crop_top = int(src_rect.get('t', '0'))
        crop_left = int(src_rect.get('l', '0'))
        crop_right = int(src_rect.get('r', '0'))

    print(f"INFO: Image crop - top={crop_top}, bottom={crop_bottom}, left={crop_left}, right={crop_right}")
    print(f"INFO: Image position - left={pic_shape.left}, top={pic_shape.top}")
    print(f"INFO: Image shape size - width={pic_shape.width}, height={pic_shape.height}")
    print(f"INFO: Slide size - width={slide_width}, height={slide_height}")

    # Component 1: Bottom crop of approximately 50% (0.4 points)
    # srcRect 'b' value should be ~50000 (50% in 1/1000th of percent units)
    # Allow tolerance of 5000 (5%)
    try:
        if crop_bottom >= 45000 and crop_bottom <= 55000:
            print(f"PASS: Component 1 - Bottom crop is {crop_bottom/1000:.1f}%, expected ~50% (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 - Bottom crop is {crop_bottom/1000:.1f}%, expected ~50%")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Visible image dimensions are approximately 3x3 inches (0.3 points)
    # The visible height = shape_height * (1 - crop_top/100000 - crop_bottom/100000)
    # The visible width = shape_width * (1 - crop_left/100000 - crop_right/100000)
    try:
        visible_height = pic_shape.height * (1.0 - crop_top / 100000.0 - crop_bottom / 100000.0)
        visible_width = pic_shape.width * (1.0 - crop_left / 100000.0 - crop_right / 100000.0)
        visible_h_inches = visible_height / 914400.0
        visible_w_inches = visible_width / 914400.0
        print(f"INFO: Visible dimensions - {visible_w_inches:.2f}\" x {visible_h_inches:.2f}\"")

        # Expect approximately 3x3 inches (allow 20% tolerance)
        width_ok = abs(visible_w_inches - 3.0) / 3.0 <= 0.20
        height_ok = abs(visible_h_inches - 3.0) / 3.0 <= 0.20
        # Also check approximate squareness (aspect ratio close to 1:1)
        if visible_w_inches > 0 and visible_h_inches > 0:
            aspect = visible_w_inches / visible_h_inches
            aspect_ok = 0.7 <= aspect <= 1.3
        else:
            aspect_ok = False

        if width_ok and height_ok and aspect_ok:
            print(f"PASS: Component 2 - Visible size ~{visible_w_inches:.2f}x{visible_h_inches:.2f} inches, near square (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 - Visible size {visible_w_inches:.2f}x{visible_h_inches:.2f} inches, expected ~3x3")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Image centered on the slide (0.3 points)
    # The visible area starts at: shape.left + shape.width * (crop_left/100000)
    # The visible area ends at: shape.left + shape.width * (1 - crop_right/100000)
    # Similarly for vertical
    try:
        vis_left = pic_shape.left + pic_shape.width * (crop_left / 100000.0)
        vis_top = pic_shape.top + pic_shape.height * (crop_top / 100000.0)
        vis_right = pic_shape.left + pic_shape.width * (1.0 - crop_right / 100000.0)
        vis_bottom = pic_shape.top + pic_shape.height * (1.0 - crop_bottom / 100000.0)

        vis_center_x = (vis_left + vis_right) / 2.0
        vis_center_y = (vis_top + vis_bottom) / 2.0
        slide_center_x = slide_width / 2.0
        slide_center_y = slide_height / 2.0

        # Tolerance: 10% of slide dimension
        x_tolerance = slide_width * 0.10
        y_tolerance = slide_height * 0.10

        x_centered = abs(vis_center_x - slide_center_x) <= x_tolerance
        y_centered = abs(vis_center_y - slide_center_y) <= y_tolerance

        print(f"INFO: Visible center = ({vis_center_x/914400:.2f}\", {vis_center_y/914400:.2f}\"), Slide center = ({slide_center_x/914400:.2f}\", {slide_center_y/914400:.2f}\")")

        if x_centered and y_centered:
            print(f"PASS: Component 3 - Image centered on slide (0.3 pts)")
            total_score += 0.3
        else:
            details = []
            if not x_centered:
                details.append(f"X off by {abs(vis_center_x - slide_center_x)/914400:.2f}\"")
            if not y_centered:
                details.append(f"Y off by {abs(vis_center_y - slide_center_y)/914400:.2f}\"")
            print(f"FAIL: Component 3 - Image not centered: {', '.join(details)}")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
