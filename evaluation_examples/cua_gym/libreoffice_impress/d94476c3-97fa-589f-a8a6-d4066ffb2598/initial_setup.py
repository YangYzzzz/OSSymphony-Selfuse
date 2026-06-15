"""
Initial Setup: Create a 7-slide botany presentation with a tall plant image on slide 4
Task ID: impress_stu_079
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_079'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
IMG_PATH = f'{WORKDIR}/plant_full.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_plant_image():
    """Create a tall plant-like image (300x600 pixels ~ 3x6 inches at 100dpi)."""
    from PIL import Image, ImageDraw, ImageFont

    width, height = 300, 600
    img = Image.new('RGB', (width, height), '#87CEEB')  # sky blue background
    draw = ImageDraw.Draw(img)

    # Ground
    draw.rectangle([0, 480, 300, 600], fill='#8B4513')  # brown soil

    # Stem
    draw.rectangle([140, 100, 160, 480], fill='#228B22')  # green stem

    # Leaves
    draw.ellipse([80, 180, 140, 260], fill='#32CD32')   # left leaf
    draw.ellipse([160, 220, 220, 300], fill='#32CD32')   # right leaf
    draw.ellipse([90, 300, 140, 370], fill='#32CD32')    # left lower leaf
    draw.ellipse([160, 340, 210, 410], fill='#32CD32')   # right lower leaf

    # Flower at top
    for angle_offset in range(0, 360, 45):
        import math
        cx, cy = 150, 90
        r = 30
        ox = cx + int(r * math.cos(math.radians(angle_offset)))
        oy = cy + int(r * math.sin(math.radians(angle_offset)))
        draw.ellipse([ox - 15, oy - 15, ox + 15, oy + 15], fill='#FF69B4')  # pink petals
    draw.ellipse([135, 75, 165, 105], fill='#FFD700')  # yellow center

    # Roots at bottom
    draw.line([150, 480, 120, 560], fill='#8B6914', width=3)
    draw.line([150, 480, 150, 570], fill='#8B6914', width=3)
    draw.line([150, 480, 180, 550], fill='#8B6914', width=3)

    img.save(IMG_PATH)
    print(f'Plant image created: {IMG_PATH}')


def create_initial():
    create_plant_image()

    prs = Presentation()
    # Standard slide size: 10 x 7.5 inches (default)

    # --- Slide 1: Title ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Introduction to Botany"
    slide1.placeholders[1].text = "Exploring the World of Plants\nDr. Elena Vasquez — Biology Department"

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Course Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Plant Cell Biology and Photosynthesis"
    for topic in [
        "Taxonomy and Classification Systems",
        "Root, Stem, and Leaf Morphology",
        "Reproduction: Pollination and Seed Dispersal",
        "Ecological Roles and Conservation",
        "Lab Practicals: Field Identification"
    ]:
        p = body2.add_paragraph()
        p.text = topic
        p.level = 0

    # --- Slide 3: Plant Cells ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Plant Cell Structure"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Key Organelles:"
    for item in [
        "Cell Wall — provides structural support",
        "Chloroplasts — site of photosynthesis",
        "Central Vacuole — stores water and nutrients",
        "Plasmodesmata — channels between adjacent cells",
        "Endoplasmic Reticulum — protein synthesis"
    ]:
        p = body3.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 4: Full Plant Diagram (with tall image) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title textbox
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(5), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Complete Plant Anatomy"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1B, 0x5E, 0x20)

    # Add the tall plant image - 3x6 inches, positioned off-center (left side)
    pic = slide4.shapes.add_picture(
        IMG_PATH,
        Inches(1.0),   # left position (not centered)
        Inches(1.0),   # top position
        Inches(3.0),   # width
        Inches(6.0),   # height
    )

    # Add description textbox on the right side
    txBox2 = slide4.shapes.add_textbox(Inches(5.0), Inches(1.5), Inches(4.5), Inches(4.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "This diagram shows the complete structure of a flowering plant, from the root system below ground to the reproductive structures at the apex."
    for run in p2.runs:
        run.font.size = Pt(14)

    p3 = tf2.add_paragraph()
    p3.text = "\nKey regions visible:"
    for run in p3.runs:
        run.font.size = Pt(14)
        run.font.bold = True

    for label in ["Flower (reproductive)", "Stem (transport)", "Leaves (photosynthesis)", "Roots (absorption)"]:
        p_item = tf2.add_paragraph()
        p_item.text = f"• {label}"
        for run in p_item.runs:
            run.font.size = Pt(12)

    # --- Slide 5: Photosynthesis ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Photosynthesis"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "6CO₂ + 6H₂O → C₆H₁₂O₆ + 6O₂"
    p5a = body5.add_paragraph()
    p5a.text = "Light-dependent reactions occur in thylakoid membranes"
    p5b = body5.add_paragraph()
    p5b.text = "Calvin cycle fixes carbon in the stroma"
    p5c = body5.add_paragraph()
    p5c.text = "Chlorophyll a absorbs red (680nm) and blue (430nm) light"

    # --- Slide 6: Pollination ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Pollination Mechanisms"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Biotic Pollination:"
    for item in ["Insects (bees, butterflies, beetles)", "Birds (hummingbirds)", "Bats (nocturnal flowers)"]:
        p = body6.add_paragraph()
        p.text = item
        p.level = 1
    p6a = body6.add_paragraph()
    p6a.text = "Abiotic Pollination:"
    for item in ["Wind (grasses, conifers)", "Water (aquatic plants)"]:
        p = body6.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 7: Summary ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Key Takeaways"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "Plants are essential for life on Earth"
    for item in [
        "Complex cellular structures enable growth and reproduction",
        "Photosynthesis converts light energy to chemical energy",
        "Diverse pollination strategies ensure species survival",
        "Understanding botany is critical for agriculture and conservation"
    ]:
        p = body7.add_paragraph()
        p.text = item

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
