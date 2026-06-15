"""
Initial Setup: Build a skill bar visualization on slide 6
Task ID: impress_design_089
Domain: libreoffice_impress

Creates an 8-slide presentation. Slide 6 has only a title 'Expertise' with no shapes.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_design_089'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Elena Moreno"
    slide1.placeholders[1].text = "Creative Designer & Visual Strategist"

    # --- Slide 2: About Me ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    tb = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "About Me"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    bio = slide2.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(4))
    tf2 = bio.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = (
        "With over 8 years of experience in visual design and brand strategy, "
        "I bring a unique blend of artistic vision and data-driven thinking to every project. "
        "My work spans digital product design, corporate identity systems, and interactive media. "
        "I have collaborated with Fortune 500 companies and startups alike, helping them craft "
        "compelling visual narratives that resonate with their audiences."
    )
    p2.runs[0].font.size = Pt(18)
    p2.runs[0].font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # --- Slide 3: Portfolio Highlights ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    tb3 = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Portfolio Highlights"
    p3.runs[0].font.size = Pt(36)
    p3.runs[0].font.bold = True
    p3.runs[0].font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    projects = [
        ("TechNova Brand Refresh", "Complete identity overhaul for a SaaS platform serving 2M+ users"),
        ("Meridian Health App", "UI/UX design for a telemedicine app with 4.8-star rating"),
        ("Flux Animation Studio", "Motion graphics package for streaming service launch campaign"),
        ("Verde Sustainable Living", "Illustration series for environmental awareness campaign"),
    ]
    y_pos = 1.8
    for title, desc in projects:
        box = slide3.shapes.add_textbox(Inches(1.5), Inches(y_pos), Inches(10), Inches(1.2))
        tframe = box.text_frame
        tframe.word_wrap = True
        pt = tframe.paragraphs[0]
        pt.text = title
        pt.runs[0].font.size = Pt(22)
        pt.runs[0].font.bold = True
        pt.runs[0].font.color.rgb = RGBColor(0x34, 0x98, 0xDB)
        pd = tframe.add_paragraph()
        pd.text = desc
        pd.runs[0].font.size = Pt(16)
        pd.runs[0].font.color.rgb = RGBColor(0x77, 0x77, 0x77)
        y_pos += 1.3

    # --- Slide 4: Experience Timeline ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    tb4 = slide4.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf4 = tb4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Experience"
    p4.runs[0].font.size = Pt(36)
    p4.runs[0].font.bold = True
    p4.runs[0].font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    roles = [
        ("2022 - Present", "Senior Creative Director", "Luminary Design Co."),
        ("2019 - 2022", "Lead Visual Designer", "PixelForge Studios"),
        ("2017 - 2019", "UI/UX Designer", "Apex Digital Agency"),
        ("2015 - 2017", "Junior Graphic Designer", "Hartwell & Partners"),
    ]
    y_pos = 1.8
    for period, role, company in roles:
        box = slide4.shapes.add_textbox(Inches(1.5), Inches(y_pos), Inches(10), Inches(1.2))
        tframe = box.text_frame
        tframe.word_wrap = True
        pt = tframe.paragraphs[0]
        pt.text = f"{period}  |  {role}"
        pt.runs[0].font.size = Pt(20)
        pt.runs[0].font.bold = True
        pt.runs[0].font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
        pc = tframe.add_paragraph()
        pc.text = company
        pc.runs[0].font.size = Pt(16)
        pc.runs[0].font.color.rgb = RGBColor(0x77, 0x77, 0x77)
        y_pos += 1.2

    # --- Slide 5: Education ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    tb5 = slide5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf5 = tb5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Education"
    p5.runs[0].font.size = Pt(36)
    p5.runs[0].font.bold = True
    p5.runs[0].font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    edu = [
        ("MFA in Visual Communication", "Rhode Island School of Design, 2015"),
        ("BFA in Graphic Design", "Savannah College of Art and Design, 2013"),
    ]
    y_pos = 1.8
    for degree, school in edu:
        box = slide5.shapes.add_textbox(Inches(1.5), Inches(y_pos), Inches(10), Inches(1))
        tframe = box.text_frame
        tframe.word_wrap = True
        pt = tframe.paragraphs[0]
        pt.text = degree
        pt.runs[0].font.size = Pt(22)
        pt.runs[0].font.bold = True
        pt.runs[0].font.color.rgb = RGBColor(0x34, 0x98, 0xDB)
        ps = tframe.add_paragraph()
        ps.text = school
        ps.runs[0].font.size = Pt(16)
        ps.runs[0].font.color.rgb = RGBColor(0x77, 0x77, 0x77)
        y_pos += 1.3

    # --- Slide 6: Expertise (TITLE ONLY - no skill bars) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    tb6 = slide6.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf6 = tb6.text_frame
    p6 = tf6.paragraphs[0]
    p6.text = "Expertise"
    p6.runs[0].font.size = Pt(36)
    p6.runs[0].font.bold = True
    p6.runs[0].font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # --- Slide 7: Testimonials ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    tb7 = slide7.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf7 = tb7.text_frame
    p7 = tf7.paragraphs[0]
    p7.text = "Testimonials"
    p7.runs[0].font.size = Pt(36)
    p7.runs[0].font.bold = True
    p7.runs[0].font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    quotes = [
        ('"Elena transformed our brand identity in ways we never imagined. '
         'Her attention to detail and creative vision are unmatched."',
         "- Sarah Patel, CEO at TechNova"),
        ('"Working with Elena on our app redesign was a game-changer. '
         'User engagement increased by 47% within the first quarter."',
         "- James Rutherford, Product Lead at Meridian Health"),
    ]
    y_pos = 1.8
    for quote, author in quotes:
        box = slide7.shapes.add_textbox(Inches(1.5), Inches(y_pos), Inches(10), Inches(1.8))
        tframe = box.text_frame
        tframe.word_wrap = True
        pq = tframe.paragraphs[0]
        pq.text = quote
        pq.runs[0].font.size = Pt(16)
        pq.runs[0].font.italic = True
        pq.runs[0].font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        pa = tframe.add_paragraph()
        pa.text = author
        pa.runs[0].font.size = Pt(14)
        pa.runs[0].font.bold = True
        pa.runs[0].font.color.rgb = RGBColor(0x34, 0x98, 0xDB)
        y_pos += 2.0

    # --- Slide 8: Contact ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    tb8 = slide8.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf8 = tb8.text_frame
    p8 = tf8.paragraphs[0]
    p8.text = "Get In Touch"
    p8.runs[0].font.size = Pt(36)
    p8.runs[0].font.bold = True
    p8.runs[0].font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    contact_info = slide8.shapes.add_textbox(Inches(3), Inches(2.5), Inches(7), Inches(3))
    ctf = contact_info.text_frame
    ctf.word_wrap = True
    lines = [
        "elena.moreno@designstudio.com",
        "+1 (415) 829-3741",
        "www.elenamoreno.design",
        "linkedin.com/in/elenamoreno",
    ]
    for i, line in enumerate(lines):
        if i == 0:
            cp = ctf.paragraphs[0]
        else:
            cp = ctf.add_paragraph()
        cp.text = line
        cp.runs[0].font.size = Pt(20)
        cp.runs[0].font.color.rgb = RGBColor(0x34, 0x98, 0xDB)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
