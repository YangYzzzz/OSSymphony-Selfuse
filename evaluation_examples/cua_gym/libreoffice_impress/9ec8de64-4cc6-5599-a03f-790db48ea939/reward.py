"""
Reward Script: Skill bar visualization on slide 6
Task ID: impress_design_089
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): Four skill labels with correct text
  Component 2 (0.30): Four background rectangles (#E0E0E0, 8x0.4in)
  Component 3 (0.30): Four overlay rectangles (#3498DB, correct widths)
  Component 4 (0.20): Correct positioning and vertical spacing
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_design_089'

# Tolerances
POS_TOL = 0.15  # inches tolerance for position
SIZE_TOL = 0.15  # inches tolerance for size

# Expected skill bars data
SKILL_BARS = [
    {"label": "Branding",     "overlay_width": 7.2, "y": 2.0},
    {"label": "UI/UX",        "overlay_width": 6.4, "y": 2.8},
    {"label": "Motion",       "overlay_width": 5.6, "y": 3.6},
    {"label": "Illustration", "overlay_width": 4.8, "y": 4.4},
]

BG_WIDTH = 8.0
BAR_HEIGHT = 0.4
BAR_X = 3.0
BG_COLOR = "E0E0E0"
OVERLAY_COLOR = "3498DB"


def emu_to_inches(emu):
    return emu / 914400.0


def approx_eq(val1, val2, tol):
    return abs(val1 - val2) <= tol


def get_fill_color_hex(shape):
    """Return hex fill color string (no #) or None."""
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # solid fill
            return str(fill.fore_color.rgb).upper()
    except Exception:
        pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide6 = prs.slides[5]  # 0-indexed
    shapes = list(slide6.shapes)

    # Classify shapes on slide 6 (excluding pre-existing title placeholder and Expertise textbox)
    labels_found = []
    bg_rects = []
    overlay_rects = []

    for shape in shapes:
        # Get position in inches
        x_in = emu_to_inches(shape.left)
        y_in = emu_to_inches(shape.top)
        w_in = emu_to_inches(shape.width)
        h_in = emu_to_inches(shape.height)

        # Check if it's a text label for a skill bar
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            for bar in SKILL_BARS:
                if text.lower() == bar["label"].lower():
                    labels_found.append({
                        "text": text,
                        "expected_label": bar["label"],
                        "y": y_in,
                        "expected_y": bar["y"],
                    })

        # Check if it's a rectangle (AUTO_SHAPE type)
        fill_hex = get_fill_color_hex(shape)
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and fill_hex is not None:
            rect_info = {
                "x": x_in, "y": y_in, "w": w_in, "h": h_in,
                "fill": fill_hex,
            }
            if fill_hex == BG_COLOR:
                bg_rects.append(rect_info)
            elif fill_hex == OVERLAY_COLOR:
                overlay_rects.append(rect_info)

    print(f"Found: {len(labels_found)} labels, {len(bg_rects)} bg rects, {len(overlay_rects)} overlay rects")

    # Component 1: Four skill labels with correct text (0.20 points)
    try:
        expected_labels = {bar["label"].lower() for bar in SKILL_BARS}
        found_labels = {l["text"].lower() for l in labels_found}
        matched_labels = expected_labels & found_labels
        label_score = len(matched_labels) / 4.0
        comp1_score = 0.20 * label_score
        if label_score == 1.0:
            print(f"PASS: Component 1 — All 4 labels found: {[l['text'] for l in labels_found]} (0.20 pts)")
            total_score += 0.20
        elif label_score > 0:
            print(f"PARTIAL: Component 1 — {len(matched_labels)}/4 labels found (matched: {matched_labels}) ({comp1_score:.2f} pts)")
            total_score += comp1_score
        else:
            print(f"FAIL: Component 1 — No skill labels found. Expected: {expected_labels}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Four background rectangles with #E0E0E0, ~8x0.4in (0.30 points)
    try:
        valid_bg = 0
        for bar in SKILL_BARS:
            found = False
            for rect in bg_rects:
                if (approx_eq(rect["y"], bar["y"], POS_TOL) and
                    approx_eq(rect["w"], BG_WIDTH, SIZE_TOL) and
                    approx_eq(rect["h"], BAR_HEIGHT, SIZE_TOL) and
                    approx_eq(rect["x"], BAR_X, POS_TOL)):
                    found = True
                    break
            if found:
                valid_bg += 1
                print(f"  BG rect for '{bar['label']}' at y={bar['y']}: OK")
            else:
                print(f"  BG rect for '{bar['label']}' at y={bar['y']}: NOT FOUND")

        bg_score = valid_bg / 4.0
        comp2_score = 0.30 * bg_score
        if bg_score == 1.0:
            print(f"PASS: Component 2 — All 4 background bars found with correct size/color (0.30 pts)")
            total_score += 0.30
        elif bg_score > 0:
            print(f"PARTIAL: Component 2 — {valid_bg}/4 background bars valid ({comp2_score:.2f} pts)")
            total_score += comp2_score
        else:
            print(f"FAIL: Component 2 — No valid background rectangles found")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Four overlay rectangles with #3498DB and correct widths (0.30 points)
    try:
        valid_overlay = 0
        for bar in SKILL_BARS:
            found = False
            for rect in overlay_rects:
                if (approx_eq(rect["y"], bar["y"], POS_TOL) and
                    approx_eq(rect["w"], bar["overlay_width"], SIZE_TOL) and
                    approx_eq(rect["h"], BAR_HEIGHT, SIZE_TOL) and
                    approx_eq(rect["x"], BAR_X, POS_TOL)):
                    found = True
                    break
            if found:
                valid_overlay += 1
                print(f"  Overlay for '{bar['label']}' width={bar['overlay_width']}in: OK")
            else:
                print(f"  Overlay for '{bar['label']}' width={bar['overlay_width']}in: NOT FOUND")

        overlay_score = valid_overlay / 4.0
        comp3_score = 0.30 * overlay_score
        if overlay_score == 1.0:
            print(f"PASS: Component 3 — All 4 overlay bars with correct widths/color (0.30 pts)")
            total_score += 0.30
        elif overlay_score > 0:
            print(f"PARTIAL: Component 3 — {valid_overlay}/4 overlay bars valid ({comp3_score:.2f} pts)")
            total_score += comp3_score
        else:
            print(f"FAIL: Component 3 — No valid overlay rectangles found")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Correct positioning — bars at x=3in, vertical spacing 0.8in (0.20 points)
    try:
        # Check that overlay rects are at x ~= 3.0 and y-positions follow 2.0, 2.8, 3.6, 4.4 spacing
        all_bar_rects = overlay_rects + bg_rects
        if len(all_bar_rects) == 0:
            print(f"FAIL: Component 4 — No bar rectangles to check positioning")
        else:
            # Sub-check A: x-position of bars ~= 3.0in (0.10 pts)
            x_correct = sum(1 for r in all_bar_rects if approx_eq(r["x"], BAR_X, POS_TOL))
            x_ratio = x_correct / len(all_bar_rects) if len(all_bar_rects) > 0 else 0
            sub_a = 0.10 * x_ratio

            # Sub-check B: y-positions match expected spacing pattern (0.10 pts)
            expected_ys = [bar["y"] for bar in SKILL_BARS]
            actual_ys = sorted(set(round(r["y"], 1) for r in overlay_rects))
            y_matches = 0
            for ey in expected_ys:
                if any(approx_eq(ay, ey, POS_TOL) for ay in actual_ys):
                    y_matches += 1
            y_ratio = y_matches / 4.0
            sub_b = 0.10 * y_ratio

            comp4_score = sub_a + sub_b
            if comp4_score >= 0.19:
                print(f"PASS: Component 4 — Positioning correct: x@3in ({x_correct}/{len(all_bar_rects)}), y-spacing ({y_matches}/4) (0.20 pts)")
                total_score += comp4_score
            elif comp4_score > 0:
                print(f"PARTIAL: Component 4 — x-pos: {x_correct}/{len(all_bar_rects)}, y-spacing: {y_matches}/4 ({comp4_score:.2f} pts)")
                total_score += comp4_score
            else:
                print(f"FAIL: Component 4 — Positioning incorrect")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
