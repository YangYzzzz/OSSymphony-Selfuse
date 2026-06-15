"""
Initial Setup: Create a 10-slide lecture presentation with slide 10 having only a title placeholder.
Task ID: impress_teach_095
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_095'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Introduction to Physical Chemistry"
    slide1.placeholders[1].text = "Prof. James Smith | Fall 2025 | Lecture 4"

    # --- Slide 2: Recap ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Recap: Chapter 3 - Chemical Kinetics"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Rate laws and reaction orders"
    body2.add_paragraph().text = "Arrhenius equation and activation energy"
    body2.add_paragraph().text = "Catalysis: homogeneous vs heterogeneous"
    body2.add_paragraph().text = "Enzyme kinetics and Michaelis-Menten model"

    # --- Slide 3: Reaction Rate Data ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Reaction Rate Constants"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Temperature dependence of rate constants:"
    items3 = [
        "k(300K) = 2.41 x 10^-3 s^-1",
        "k(350K) = 8.67 x 10^-2 s^-1",
        "k(400K) = 1.54 x 10^0 s^-1",
        "Ea = 75.3 kJ/mol (from Arrhenius plot)",
    ]
    for item in items3:
        body3.add_paragraph().text = item

    # --- Slide 4: Equilibrium Concepts ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Chemical Equilibrium Fundamentals"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Le Chatelier's Principle"
    body4.add_paragraph().text = "Equilibrium constant expressions (Kc, Kp)"
    body4.add_paragraph().text = "Relationship between Gibbs free energy and K"
    body4.add_paragraph().text = "Delta_G = -RT ln(K)"

    # --- Slide 5: Phase Diagrams ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Phase Diagrams and Phase Transitions"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Triple point of water: 273.16 K, 611.73 Pa"
    body5.add_paragraph().text = "Critical point of CO2: 304.25 K, 7.39 MPa"
    body5.add_paragraph().text = "Clausius-Clapeyron equation applications"
    body5.add_paragraph().text = "Supercritical fluid extraction in industry"

    # --- Slide 6: Thermodynamic Laws ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Laws of Thermodynamics - Overview"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "0th Law: Thermal equilibrium is transitive"
    body6.add_paragraph().text = "1st Law: Energy conservation (Delta_U = q + w)"
    body6.add_paragraph().text = "2nd Law: Entropy of isolated system always increases"
    body6.add_paragraph().text = "3rd Law: Entropy approaches zero as T approaches 0 K"

    # --- Slide 7: Enthalpy Calculations ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Enthalpy of Formation - Selected Values"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "H2O(l): -285.8 kJ/mol"
    body7.add_paragraph().text = "CO2(g): -393.5 kJ/mol"
    body7.add_paragraph().text = "NaCl(s): -411.2 kJ/mol"
    body7.add_paragraph().text = "CH4(g): -74.8 kJ/mol"
    body7.add_paragraph().text = "C2H5OH(l): -277.7 kJ/mol"

    # --- Slide 8: Electrochemistry ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Electrochemistry Basics"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Standard electrode potentials (E0)"
    body8.add_paragraph().text = "Nernst equation: E = E0 - (RT/nF)ln(Q)"
    body8.add_paragraph().text = "Galvanic cells vs electrolytic cells"
    body8.add_paragraph().text = "Faraday's laws of electrolysis"

    # --- Slide 9: Practice Problems ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Practice Problems"
    body9 = slide9.placeholders[1].text_frame
    body9.text = "1. Calculate Delta_H for the combustion of methane."
    body9.add_paragraph().text = "2. Determine K at 500K given K(300K) and Delta_H."
    body9.add_paragraph().text = "3. Sketch the phase diagram for a substance with these properties."
    body9.add_paragraph().text = "4. Calculate E_cell for Zn/Cu galvanic cell at non-standard conditions."

    # --- Slide 10: End slide (title placeholder only, no custom content) ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add only a simple title textbox at the top
    txBox = slide10.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "End of Lecture 4"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
