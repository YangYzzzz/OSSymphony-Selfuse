"""
Reward Script: End-of-lecture slide with Questions?, decorative line, next class info, and email
Task ID: impress_teach_095
Domain: libreoffice_impress
Scoring:
  Component 1 (0.35): "Questions?" text - centered, 54pt, bold, #1A237E
  Component 2 (0.25): Gold decorative horizontal line (#FFD700, ~3pt height)
  Component 3 (0.20): "Next class: Chapter 5 - Thermodynamics" text - centered, 20pt
  Component 4 (0.20): "Email: prof.smith@university.edu" text - centered, 16pt, #757575
"""

import os
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_095'


def find_text_shapes(slide):
    """Get all shapes with text frames on a slide."""
    results = []
    for shape in slide.shapes:
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            results.append(shape)
    return results


def find_shape_by_text(slide, target_text):
    """Find a shape whose text content contains target_text."""
    target_lower = target_text.lower().strip()
    for shape in slide.shapes:
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            full_text = shape.text_frame.text.strip().lower()
            if target_lower in full_text:
                return shape
    return None


def find_line_shape(slide):
    """Find a shape that acts as a decorative horizontal line (thin rectangle or line)."""
    for shape in slide.shapes:
        # A decorative line is typically a very thin rectangle (height <= ~5pt = 63500 EMU)
        # or an actual line connector
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # Thin rectangle used as line
            if shape.height <= Pt(10) and shape.width > shape.height * 5:
                return shape
        # Freeform or line shapes
        if shape.shape_type in (MSO_SHAPE_TYPE.FREEFORM, ):
            if shape.height <= Pt(10) and shape.width > shape.height * 3:
                return shape
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 10 slides
    if len(prs.slides) < 10:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 10")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[9]  # Slide 10 (0-indexed)

    # Component 1: "Questions?" text with correct formatting (0.35 points)
    try:
        questions_shape = find_shape_by_text(slide, "Questions?")
        if questions_shape is not None:
            para = None
            target_run = None
            for p in questions_shape.text_frame.paragraphs:
                for r in p.runs:
                    if "questions?" in (r.text or "").lower().strip():
                        para = p
                        target_run = r
                        break
                if target_run:
                    break

            if target_run is None:
                print("FAIL: Component 1 -- 'Questions?' run not found")
            else:
                sub_score = 0.0
                # Check text content
                if "questions?" in target_run.text.lower().strip():
                    sub_score += 0.05
                    print(f"PASS: Component 1a -- Text is '{target_run.text}'")
                else:
                    print(f"FAIL: Component 1a -- Expected 'Questions?', got '{target_run.text}'")

                # Check centered alignment
                align = para.alignment
                if align == PP_ALIGN.CENTER:
                    sub_score += 0.05
                    print("PASS: Component 1b -- Text is centered")
                else:
                    print(f"FAIL: Component 1b -- Expected CENTER alignment, got {align}")

                # Check font size ~54pt (685800 EMU)
                expected_size = Pt(54)
                actual_size = target_run.font.size
                if actual_size is not None and abs(actual_size - expected_size) <= Pt(2):
                    sub_score += 0.10
                    print(f"PASS: Component 1c -- Font size {actual_size/12700:.1f}pt (expected ~54pt)")
                else:
                    print(f"FAIL: Component 1c -- Expected ~54pt, got {actual_size/12700:.1f}pt" if actual_size else "FAIL: Component 1c -- Font size is None")

                # Check bold
                if target_run.font.bold is True:
                    sub_score += 0.05
                    print("PASS: Component 1d -- Text is bold")
                else:
                    print(f"FAIL: Component 1d -- Expected bold=True, got {target_run.font.bold}")

                # Check color #1A237E (dark blue)
                try:
                    if target_run.font.color.type is not None:
                        rgb = str(target_run.font.color.rgb).upper()
                        if rgb == "1A237E":
                            sub_score += 0.10
                            print(f"PASS: Component 1e -- Color is #1A237E")
                        else:
                            print(f"FAIL: Component 1e -- Expected #1A237E, got #{rgb}")
                    else:
                        print("FAIL: Component 1e -- No explicit color set")
                except Exception as e:
                    print(f"FAIL: Component 1e -- Color check error: {e}")

                if sub_score > 0:
                    total_score += sub_score
                print(f"  Component 1 subtotal: {sub_score}/0.35")
        else:
            print("FAIL: Component 1 -- No shape with 'Questions?' text found on slide 10")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Gold decorative horizontal line (0.25 points)
    try:
        line_shape = find_line_shape(slide)
        if line_shape is not None:
            sub_score = 0.0

            # Check that it's thin (acts as a line) - height should be ~3pt (38100 EMU)
            height_pt = line_shape.height / 12700
            if height_pt <= 6:
                sub_score += 0.10
                print(f"PASS: Component 2a -- Line height {height_pt:.1f}pt (thin rectangle as line)")
            else:
                print(f"FAIL: Component 2a -- Expected thin line, height is {height_pt:.1f}pt")

            # Check gold color (#FFD700) in fill
            try:
                fill = line_shape.fill
                if fill.type is not None:
                    fore_rgb = str(fill.fore_color.rgb).upper()
                    if fore_rgb == "FFD700":
                        sub_score += 0.15
                        print(f"PASS: Component 2b -- Fill color is gold (#FFD700)")
                    else:
                        print(f"FAIL: Component 2b -- Expected gold #FFD700, got #{fore_rgb}")
                else:
                    print("FAIL: Component 2b -- No fill color set on line shape")
            except Exception as e:
                print(f"FAIL: Component 2b -- Fill color check error: {e}")

            if sub_score > 0:
                total_score += sub_score
            print(f"  Component 2 subtotal: {sub_score}/0.25")
        else:
            print("FAIL: Component 2 -- No decorative line shape found on slide 10")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: "Next class: Chapter 5 - Thermodynamics" text, centered, 20pt (0.20 points)
    try:
        next_class_shape = find_shape_by_text(slide, "Next class")
        if next_class_shape is not None:
            para = next_class_shape.text_frame.paragraphs[0]
            runs = [r for r in para.runs if (r.text or "").strip()]
            sub_score = 0.0

            # Check text content
            full_text = para.text.strip()
            if "next class" in full_text.lower() and "chapter 5" in full_text.lower() and "thermodynamics" in full_text.lower():
                sub_score += 0.08
                print(f"PASS: Component 3a -- Text: '{full_text}'")
            else:
                print(f"FAIL: Component 3a -- Expected 'Next class: Chapter 5 - Thermodynamics', got '{full_text}'")

            # Check centered
            if para.alignment == PP_ALIGN.CENTER:
                sub_score += 0.04
                print("PASS: Component 3b -- Text is centered")
            else:
                print(f"FAIL: Component 3b -- Expected CENTER, got {para.alignment}")

            # Check font size ~20pt
            if runs:
                actual_size = runs[0].font.size
                if actual_size is not None and abs(actual_size - Pt(20)) <= Pt(2):
                    sub_score += 0.08
                    print(f"PASS: Component 3c -- Font size {actual_size/12700:.1f}pt (expected ~20pt)")
                else:
                    size_str = f"{actual_size/12700:.1f}pt" if actual_size else "None"
                    print(f"FAIL: Component 3c -- Expected ~20pt, got {size_str}")
            else:
                print("FAIL: Component 3c -- No runs found")

            if sub_score > 0:
                total_score += sub_score
            print(f"  Component 3 subtotal: {sub_score}/0.20")
        else:
            print("FAIL: Component 3 -- No shape with 'Next class' text found on slide 10")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Email text, centered, 16pt, #757575 (0.20 points)
    try:
        email_shape = find_shape_by_text(slide, "prof.smith@university.edu")
        if email_shape is not None:
            para = email_shape.text_frame.paragraphs[0]
            runs = [r for r in para.runs if (r.text or "").strip()]
            sub_score = 0.0

            # Check text content
            full_text = para.text.strip()
            if "prof.smith@university.edu" in full_text.lower():
                sub_score += 0.05
                print(f"PASS: Component 4a -- Text: '{full_text}'")
            else:
                print(f"FAIL: Component 4a -- Expected email text, got '{full_text}'")

            # Check centered
            if para.alignment == PP_ALIGN.CENTER:
                sub_score += 0.03
                print("PASS: Component 4b -- Text is centered")
            else:
                print(f"FAIL: Component 4b -- Expected CENTER, got {para.alignment}")

            # Check font size ~16pt
            if runs:
                actual_size = runs[0].font.size
                if actual_size is not None and abs(actual_size - Pt(16)) <= Pt(2):
                    sub_score += 0.05
                    print(f"PASS: Component 4c -- Font size {actual_size/12700:.1f}pt (expected ~16pt)")
                else:
                    size_str = f"{actual_size/12700:.1f}pt" if actual_size else "None"
                    print(f"FAIL: Component 4c -- Expected ~16pt, got {size_str}")
            else:
                print("FAIL: Component 4c -- No runs found")

            # Check color #757575
            if runs:
                try:
                    if runs[0].font.color.type is not None:
                        rgb = str(runs[0].font.color.rgb).upper()
                        if rgb == "757575":
                            sub_score += 0.07
                            print(f"PASS: Component 4d -- Color is #757575")
                        else:
                            print(f"FAIL: Component 4d -- Expected #757575, got #{rgb}")
                    else:
                        print("FAIL: Component 4d -- No explicit color set")
                except Exception as e:
                    print(f"FAIL: Component 4d -- Color check error: {e}")
            else:
                print("FAIL: Component 4d -- No runs to check color")

            if sub_score > 0:
                total_score += sub_score
            print(f"  Component 4 subtotal: {sub_score}/0.20")
        else:
            print("FAIL: Component 4 -- No shape with email text found on slide 10")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
