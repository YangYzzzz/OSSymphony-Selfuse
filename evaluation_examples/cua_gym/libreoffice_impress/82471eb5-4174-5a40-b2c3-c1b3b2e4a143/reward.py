"""
Reward Script: Product Catalog Presentation
Task ID: impress_wf_031
Domain: libreoffice_impress
Scoring:
  C1 (0.15): 7 slides total
  C2 (0.10): Slide 1 title "Product Catalog 2024"
  C3 (0.15): Slide 2 has 4 category buttons with correct text
  C4 (0.15): Slide 2 buttons have hyperlinks (action jump) to slides 3-6
  C5 (0.15): Slides 3-6 each have 4 product cards with text
  C6 (0.10): Slides 3-6 each have "Back to Categories" hyperlink
  C7 (0.10): Slide 7 has a table
  C8 (0.10): Slide 2 buttons have shadow and border
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_031'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.oxml.ns import qn
    except ImportError as e:
        print(f"CRITICAL: Missing library: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Presentation has exactly 7 slides (0.15 points)
    try:
        if num_slides == 7:
            print(f"PASS: Component 1 — slide count is 7 (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — expected 7 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Gate: need at least 7 slides to check further
    if num_slides < 7:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Slide 1 title contains "Product Catalog 2024" (0.10 points)
    try:
        slide1 = prs.slides[0]
        slide1_text = " ".join(
            shape.text for shape in slide1.shapes if hasattr(shape, 'text') and shape.text
        ).strip()
        if "Product Catalog 2024" in slide1_text:
            print(f"PASS: Component 2 — Slide 1 contains 'Product Catalog 2024' (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 — Slide 1 text: '{slide1_text[:100]}', expected 'Product Catalog 2024'")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 2 has 4 category buttons (auto shapes) with correct text (0.15 points)
    try:
        slide2 = prs.slides[1]
        expected_categories = {"electronics", "clothing", "home", "sports"}
        found_categories = set()
        button_shapes = []
        for shape in slide2.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                text_lower = shape.text.strip().lower()
                if text_lower in expected_categories:
                    found_categories.add(text_lower)
                    button_shapes.append(shape)

        if found_categories == expected_categories:
            print(f"PASS: Component 3 — found all 4 category buttons: {found_categories} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 3 — found categories: {found_categories}, missing: {expected_categories - found_categories}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 2 buttons have hyperlinks (slide jump actions) to slides 3-6 (0.15 points)
    try:
        slide2 = prs.slides[1]
        buttons_with_links = 0
        for shape in slide2.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                el = shape._element
                # Check for hlinkClick on cNvPr (shape-level hyperlink / action)
                cNvPr = el.find('.//' + qn('p:cNvPr'))
                has_link = False
                if cNvPr is not None:
                    hlinkClick = cNvPr.find(qn('a:hlinkClick'))
                    if hlinkClick is not None:
                        action = hlinkClick.get('action', '')
                        if 'hlinksldjump' in action.lower():
                            has_link = True

                # Also check run-level hyperlinks
                if not has_link and shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.hyperlink and run.hyperlink.address:
                                if 'slide' in str(run.hyperlink.address).lower():
                                    has_link = True
                                    break

                if has_link:
                    buttons_with_links += 1

        if buttons_with_links >= 4:
            print(f"PASS: Component 4 — {buttons_with_links} buttons have slide hyperlinks (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 — only {buttons_with_links} buttons have slide hyperlinks, expected 4")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slides 3-6 each have 4 product cards (auto shapes with text) (0.15 points)
    try:
        category_names = ["Electronics", "Clothing", "Home", "Sports"]
        slides_with_products = 0
        for si in range(2, 6):  # slides 3-6 (0-indexed 2-5)
            slide = prs.slides[si]
            product_cards = 0
            for shape in slide.shapes:
                # Product cards are auto shapes with product text (name + price)
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.has_text_frame:
                    text = shape.text.strip()
                    # Product cards have price indicators ($, or multi-line with product info)
                    if '$' in text or len(text.split('\n')) >= 2:
                        product_cards += 1
            if product_cards >= 4:
                slides_with_products += 1
                print(f"  Slide {si+1}: {product_cards} product cards found")
            else:
                print(f"  Slide {si+1}: only {product_cards} product cards found")

        if slides_with_products == 4:
            print(f"PASS: Component 5 — all 4 category slides have 4+ product cards (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 5 — {slides_with_products}/4 category slides have 4+ product cards")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Slides 3-6 each have "Back to Categories" hyperlink (0.10 points)
    try:
        slides_with_back_link = 0
        for si in range(2, 6):  # slides 3-6
            slide = prs.slides[si]
            has_back_link = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    full_text = shape.text.strip().lower()
                    if "back to categories" in full_text:
                        # Check for hyperlink
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.hyperlink and run.hyperlink.address:
                                    has_back_link = True
                                    break
                            if has_back_link:
                                break
                    # Also check shape-level hyperlink
                    if not has_back_link and "back to categories" in full_text:
                        el = shape._element
                        cNvPr = el.find('.//' + qn('p:cNvPr'))
                        if cNvPr is not None:
                            hlinkClick = cNvPr.find(qn('a:hlinkClick'))
                            if hlinkClick is not None:
                                has_back_link = True
                if has_back_link:
                    break

            if has_back_link:
                slides_with_back_link += 1
                print(f"  Slide {si+1}: has 'Back to Categories' hyperlink")
            else:
                print(f"  Slide {si+1}: missing 'Back to Categories' hyperlink")

        if slides_with_back_link == 4:
            print(f"PASS: Component 6 — all 4 category slides have back link (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 6 — {slides_with_back_link}/4 category slides have back link")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Slide 7 has a table (0.10 points)
    try:
        slide7 = prs.slides[6]
        has_table = False
        for shape in slide7.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                if len(table.rows) >= 2 and len(table.columns) >= 2:
                    has_table = True
                    print(f"  Table found: {len(table.rows)} rows x {len(table.columns)} cols")
                    break

        if has_table:
            print(f"PASS: Component 7 — Slide 7 has a table (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 7 — Slide 7 has no table with 2+ rows/cols")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    # Component 8: Slide 2 buttons have shadow effects AND visible borders (0.10 points)
    try:
        slide2 = prs.slides[1]
        buttons_with_effects = 0
        for shape in slide2.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                el = shape._element
                # spPr is under p: namespace (p:spPr), find it via either namespace
                sp_pr = el.find(qn('p:spPr'))
                if sp_pr is None:
                    sp_pr = el.find('.//' + qn('a:spPr'))
                if sp_pr is None:
                    continue

                # Check shadow in effectLst
                has_shadow = False
                effect_lst = sp_pr.find(qn('a:effectLst'))
                if effect_lst is not None:
                    outer_shdw = effect_lst.find(qn('a:outerShdw'))
                    if outer_shdw is not None:
                        has_shadow = True

                # Check border (a:ln with width > 0)
                has_border = False
                ln = sp_pr.find(qn('a:ln'))
                if ln is not None:
                    w = ln.get('w')
                    if w is not None and int(w) > 0:
                        has_border = True

                if has_shadow and has_border:
                    buttons_with_effects += 1

        if buttons_with_effects >= 4:
            print(f"PASS: Component 8 — {buttons_with_effects} buttons have shadow+border (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 8 — only {buttons_with_effects} buttons have shadow+border, expected 4")
    except Exception as e:
        print(f"ERROR: Component 8 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = os.path.join(WORKDIR, TASK_ID + '.pptx')
if not os.path.exists(file_path):
    # Also check Desktop
    alt_path = os.path.join(WORKDIR, 'Desktop', 'Product_Catalog.pptx')
    if os.path.exists(alt_path):
        file_path = alt_path
    else:
        print(f"File not found: {file_path}")
        print("REWARD: 0.0")
        import sys
        sys.exit(0)

verify_task(file_path)
