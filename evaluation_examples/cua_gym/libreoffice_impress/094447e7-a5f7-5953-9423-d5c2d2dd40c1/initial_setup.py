"""
Initial Setup: Create an 8-slide Product Launch presentation.
Slide 3 has a blue background rectangle, title 'Key Features', and 5 bullet points.
No animations are applied to any objects.
Task ID: impress_gf2_001
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_001'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    slide_width = prs.slide_width   # default 10 inches
    slide_height = prs.slide_height  # default 7.5 inches

    # ============================================================
    # Slide 1: Title Slide — "Product Launch 2024"
    # ============================================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Product Launch 2024"
    slide1.placeholders[1].text = "Revolutionizing the Market"

    # ============================================================
    # Slide 2: "Market Overview"
    # ============================================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Market Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "The global market has seen a 23% increase in demand for smart home devices."
    p2 = body2.add_paragraph()
    p2.text = "Our target segment grew from $4.2B to $5.8B in 2023."
    p3 = body2.add_paragraph()
    p3.text = "Consumer surveys show 78% interest in our product category."

    # ============================================================
    # Slide 3: "Key Features" — The target slide for animations
    # Uses BLANK layout. All shapes created manually.
    # ============================================================
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only layout

    # 3a. Blue background rectangle (full-slide size)
    rect = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height
    )
    rect.name = "BackgroundRect"
    fill = rect.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x00, 0x3C, 0x71)  # Dark blue
    rect.line.fill.background()  # No border

    # 3b. Title text box "Key Features"
    title_box = slide3.shapes.add_textbox(
        Inches(1.5), Inches(0.5), Inches(7), Inches(1)
    )
    title_box.name = "TitleBox"
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = "Key Features"
    p_title.alignment = PP_ALIGN.CENTER
    run_title = p_title.runs[0]
    run_title.font.name = "Arial"
    run_title.font.size = Pt(36)
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 3c. Bullet list with 5 items
    bullet_box = slide3.shapes.add_textbox(
        Inches(1.5), Inches(2.0), Inches(7), Inches(4.5)
    )
    bullet_box.name = "BulletList"
    tf_bullets = bullet_box.text_frame
    tf_bullets.word_wrap = True

    bullets = [
        "Advanced AI-powered voice recognition with 99.2% accuracy",
        "Seamless integration with over 200 smart home devices",
        "Ultra-low power consumption — 18 months on a single charge",
        "Military-grade AES-256 encryption for all communications",
        "Customizable routines via intuitive drag-and-drop interface",
    ]

    for i, text in enumerate(bullets):
        if i == 0:
            p = tf_bullets.paragraphs[0]
        else:
            p = tf_bullets.add_paragraph()
        p.text = text
        p.level = 0
        p.space_after = Pt(12)
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # ============================================================
    # Slide 4: "Technical Specifications"
    # ============================================================
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Technical Specifications"
    body4 = slide4.placeholders[1].text_frame
    specs = [
        "Processor: Qualcomm Snapdragon 8 Gen 3",
        "RAM: 8GB LPDDR5X",
        "Connectivity: Wi-Fi 7, Bluetooth 5.4, Thread, Matter",
        "Dimensions: 142mm × 68mm × 12mm",
        "Weight: 185g",
    ]
    body4.text = specs[0]
    for s in specs[1:]:
        body4.add_paragraph().text = s

    # ============================================================
    # Slide 5: "Competitive Analysis"
    # ============================================================
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Competitive Analysis"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Our product outperforms competitors in 4 of 5 key categories."
    body5.add_paragraph().text = "Price point is 15% below the market average."
    body5.add_paragraph().text = "Battery life exceeds nearest competitor by 6 months."

    # ============================================================
    # Slide 6: "Go-to-Market Strategy"
    # ============================================================
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Go-to-Market Strategy"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Phase 1 (Q1 2024): Beta launch with 1,000 early adopters"
    body6.add_paragraph().text = "Phase 2 (Q2 2024): Full retail rollout across North America"
    body6.add_paragraph().text = "Phase 3 (Q3 2024): European and Asian market expansion"
    body6.add_paragraph().text = "Phase 4 (Q4 2024): Enterprise and B2B partnerships"

    # ============================================================
    # Slide 7: "Financial Projections"
    # ============================================================
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Financial Projections"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "Year 1 Revenue: $12.5M (conservative estimate)"
    body7.add_paragraph().text = "Year 2 Revenue: $34.8M (with enterprise channel)"
    body7.add_paragraph().text = "Break-even point: Month 14"
    body7.add_paragraph().text = "Gross margin target: 62%"

    # ============================================================
    # Slide 8: "Thank You / Q&A"
    # ============================================================
    slide8 = prs.slides.add_slide(prs.slide_layouts[0])
    slide8.shapes.title.text = "Thank You"
    slide8.placeholders[1].text = "Questions & Answers"

    # Save
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
