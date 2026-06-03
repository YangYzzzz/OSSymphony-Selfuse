"""
Initial Setup: Budget Report presentation with 6 slides, table on slide 3 with default borders
Task ID: impress_tct_003
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_003'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_cell_text(cell, text, font_name="Calibri", font_size=Pt(11),
                  bold=False, color=None, alignment=None):
    """Helper to set cell text with formatting."""
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if alignment:
        p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def set_thin_gray_borders(table):
    """Set thin (0.5pt) gray borders on all cells — default-looking borders."""
    tbl = table._tbl
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    for row_idx in range(len(table.rows)):
        for col_idx in range(len(table.columns)):
            tc = tbl.findall(qn('a:tr'))[row_idx].findall(qn('a:tc'))[col_idx]
            tcPr = tc.find(qn('a:tcPr'))
            if tcPr is None:
                tcPr = tc.makeelement(qn('a:tcPr'), {})
                tc.insert(0, tcPr)

            # Remove existing borders
            for border_tag in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                existing = tcPr.find(qn(border_tag))
                if existing is not None:
                    tcPr.remove(existing)

            # Add thin gray borders (0.5pt = 6350 EMU)
            for border_tag in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                ln = tcPr.makeelement(qn(border_tag), {'w': '6350', 'cmpd': 'sng'})
                solidFill = ln.makeelement(qn('a:solidFill'), {})
                srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': 'BFBFBF'})
                solidFill.append(srgbClr)
                ln.append(solidFill)
                prstDash = ln.makeelement(qn('a:prstDash'), {'val': 'solid'})
                ln.append(prstDash)
                tcPr.append(ln)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q4 2025 Budget Report"
    slide1.placeholders[1].text = "Finance Department — Prepared by Rachel Torres"

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Total annual expenditure: $2.87M across 4 departments"
    p2a = tf2.add_paragraph()
    p2a.text = "Engineering remains the largest cost center at 42% of total budget"
    p2a.level = 0
    p2b = tf2.add_paragraph()
    p2b.text = "Marketing spend increased 15% QoQ driven by Q4 campaigns"
    p2b.level = 0
    p2c = tf2.add_paragraph()
    p2c.text = "Operations achieved 8% cost reduction through process optimization"
    p2c.level = 0

    # --- Slide 3: Budget Breakdown Table ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title textbox
    title_box = slide3.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(10), Inches(0.8))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.alignment = PP_ALIGN.LEFT
    run_title = p_title.add_run()
    run_title.text = "Departmental Budget Breakdown (in $thousands)"
    run_title.font.name = "Calibri"
    run_title.font.size = Pt(24)
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # 4 rows x 5 columns table
    rows, cols = 4, 5
    table_shape = slide3.shapes.add_table(
        rows, cols,
        Inches(1.0), Inches(1.5), Inches(10.5), Inches(3.5)
    )
    table = table_shape.table

    # Set column widths
    table.columns[0].width = Inches(2.5)
    for c in range(1, 5):
        table.columns[c].width = Inches(2.0)

    # Headers
    headers = ["Department", "Q1 Budget", "Q2 Budget", "Q3 Budget", "Q4 Budget"]
    header_color = RGBColor(0x33, 0x33, 0x33)
    for c, h in enumerate(headers):
        set_cell_text(table.cell(0, c), h, bold=True, color=header_color,
                      font_size=Pt(12), alignment=PP_ALIGN.CENTER)

    # Data rows
    data = [
        ["Engineering", "$312,500", "$298,400", "$305,200", "$289,100"],
        ["Marketing", "$145,800", "$152,300", "$168,900", "$175,600"],
        ["Operations", "$98,200", "$95,700", "$91,400", "$88,300"],
    ]
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            align = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            set_cell_text(table.cell(r, c), val, font_size=Pt(11), alignment=align)

    # Set thin gray borders (default looking)
    set_thin_gray_borders(table)

    # --- Slide 4: Trend Analysis ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Year-over-Year Trend Analysis"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Engineering: Stable with 3% reduction in Q4 due to headcount optimization"
    p4a = tf4.add_paragraph()
    p4a.text = "Marketing: 20% increase year-over-year, driven by digital campaigns"
    p4b = tf4.add_paragraph()
    p4b.text = "Operations: Consistent downward trend (-8% YoY) from automation initiatives"
    p4c = tf4.add_paragraph()
    p4c.text = "R&D allocation pending board approval for FY2026"

    # --- Slide 5: Key Metrics ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title_box5 = slide5.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(10), Inches(0.8))
    tf5 = title_box5.text_frame
    p5 = tf5.paragraphs[0]
    run5 = p5.add_run()
    run5.text = "Key Financial Metrics"
    run5.font.name = "Calibri"
    run5.font.size = Pt(24)
    run5.font.bold = True
    run5.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    metrics = [
        ("Total Revenue", "$4.23M", "+12% YoY"),
        ("Total Expenses", "$2.87M", "+5% YoY"),
        ("Net Margin", "32.2%", "+3.1pp"),
        ("Cost per Employee", "$127,400", "-2% YoY"),
    ]
    for i, (label, value, change) in enumerate(metrics):
        box = slide5.shapes.add_textbox(
            Inches(1.0 + (i % 2) * 5.5), Inches(1.5 + (i // 2) * 2.5),
            Inches(4.5), Inches(2.0)
        )
        tf = box.text_frame
        p_label = tf.paragraphs[0]
        r_label = p_label.add_run()
        r_label.text = label
        r_label.font.size = Pt(14)
        r_label.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

        p_val = tf.add_paragraph()
        r_val = p_val.add_run()
        r_val.text = value
        r_val.font.size = Pt(28)
        r_val.font.bold = True
        r_val.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

        p_chg = tf.add_paragraph()
        r_chg = p_chg.add_run()
        r_chg.text = change
        r_chg.font.size = Pt(12)
        r_chg.font.color.rgb = RGBColor(0x2E, 0x7D, 0x32)

    # --- Slide 6: Next Steps ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Next Steps & Recommendations"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Finalize Q1 2026 budget allocations by January 15th"
    p6a = tf6.add_paragraph()
    p6a.text = "Schedule departmental budget review meetings for week of Jan 6"
    p6b = tf6.add_paragraph()
    p6b.text = "Submit capital expenditure proposals for board review"
    p6c = tf6.add_paragraph()
    p6c.text = "Implement new expense tracking system by end of Q1 2026"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
