"""
Initial Setup: Research Symposium Presentation - Multi-Monitor Configuration
Task ID: impress_gf5_041
Domain: libreoffice_impress

Creates a 25-slide research symposium presentation with speaker notes.
No presenter console or timer settings configured.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_041'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_placeholder(placeholder, text, font_name="Calibri", font_size=18, bold=False, color=None):
    """Helper to set text with formatting on a placeholder."""
    tf = placeholder.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    for run in p.runs:
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color


def add_bullet_content(text_frame, items, font_name="Calibri", font_size=16):
    """Add bulleted list to a text frame."""
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide data: (title, subtitle/content, notes)
    slides_data = [
        # Slide 1: Title
        {
            "layout": 0,
            "title": "Advances in Neural Architecture Search for Edge Computing",
            "subtitle": "12th International Research Symposium on Computational Intelligence\nDr. Amara Okafor, Prof. Hiroshi Tanaka, Dr. Elena Vasquez\nMarch 28, 2025 — Singapore",
            "notes": "Welcome everyone. Today we present our findings on NAS for edge devices. Acknowledge funding from NSF grant #2024-1847 and Horizon Europe."
        },
        # Slide 2: Agenda
        {
            "layout": 1,
            "title": "Symposium Agenda",
            "content": ["Research motivation and problem statement", "Literature review: NAS methods 2020-2025",
                        "Proposed framework: EdgeNAS-Pro", "Experimental methodology",
                        "Results on benchmark datasets", "Ablation studies",
                        "Real-world deployment case studies", "Future directions and Q&A"],
            "notes": "Brief overview of our 45-minute session. We'll leave 10 minutes for questions at the end."
        },
        # Slide 3: Problem Statement
        {
            "layout": 1,
            "title": "Problem Statement",
            "content": ["Edge devices have strict latency constraints (<10ms inference)",
                        "Manual architecture design is time-consuming and sub-optimal",
                        "Existing NAS methods require 100+ GPU hours",
                        "Gap: No efficient NAS framework targeting heterogeneous edge hardware"],
            "notes": "Emphasize the practical gap. Show statistics: 73% of ML engineers spend >2 weeks on architecture tuning per project."
        },
        # Slide 4: Research Questions
        {
            "layout": 1,
            "title": "Research Questions",
            "content": ["RQ1: Can we reduce NAS search cost below 5 GPU hours for edge targets?",
                        "RQ2: Does hardware-aware search improve latency-accuracy tradeoff?",
                        "RQ3: How does EdgeNAS-Pro generalize across ARM, RISC-V, and FPGA?",
                        "RQ4: What is the impact on real-world deployment metrics?"],
            "notes": "These four questions structure our entire presentation. RQ1 is our primary contribution."
        },
        # Slide 5: Literature Review
        {
            "layout": 1,
            "title": "Literature Review: NAS Landscape",
            "content": ["EfficientNet (Tan & Le, 2019): Compound scaling",
                        "Once-for-All (Cai et al., 2020): Progressive shrinking",
                        "FBNetV3 (Dai et al., 2021): Accuracy predictor approach",
                        "DONNA (Moons et al., 2022): Multi-objective optimization",
                        "AutoFormer (Chen et al., 2023): Transformer NAS"],
            "notes": "Key insight: all prior methods focus on single hardware target. Our contribution is the multi-target approach."
        },
        # Slide 6: Research Gap
        {
            "layout": 1,
            "title": "Identified Research Gap",
            "content": ["No unified search space for heterogeneous edge hardware",
                        "Transfer learning between hardware targets unexplored",
                        "Deployment-aware metrics missing from search objectives",
                        "Power consumption rarely considered in NAS formulation"],
            "notes": "This slide is critical for establishing novelty. Pause here for emphasis."
        },
        # Slide 7: Proposed Framework
        {
            "layout": 1,
            "title": "EdgeNAS-Pro Framework Overview",
            "content": ["Three-phase progressive search strategy",
                        "Hardware-agnostic supernet with device-specific adapters",
                        "Pareto-optimal architecture selection engine",
                        "Automated deployment pipeline with profiling"],
            "notes": "Show the architecture diagram. Explain each phase: macro search, micro search, deployment optimization."
        },
        # Slide 8: Phase 1 - Macro Search
        {
            "layout": 1,
            "title": "Phase 1: Macro Architecture Search",
            "content": ["Search space: depth, width multipliers, resolution",
                        "Evolutionary algorithm with 500 candidates per generation",
                        "Fitness function: accuracy + latency Pareto rank",
                        "Convergence in 2.3 GPU hours (avg over 50 trials)"],
            "notes": "This phase reduces search space by 94%. Show convergence plot."
        },
        # Slide 9: Phase 2 - Micro Search
        {
            "layout": 1,
            "title": "Phase 2: Micro Operation Search",
            "content": ["Operation set: MBConv, FusedMBConv, Attention blocks",
                        "Differentiable search with Gumbel-Softmax relaxation",
                        "Hardware-aware latency lookup tables",
                        "Joint optimization across 3 target platforms"],
            "notes": "Micro search takes 1.8 GPU hours. Total framework cost: 4.1 hours, well under our 5-hour target."
        },
        # Slide 10: Phase 3 - Deployment
        {
            "layout": 1,
            "title": "Phase 3: Deployment Optimization",
            "content": ["Post-training quantization: INT8 and mixed-precision",
                        "Layer fusion and memory layout optimization",
                        "Runtime profiling with 95th percentile latency target",
                        "Automated model packaging for TFLite, ONNX, TVM"],
            "notes": "Deployment phase adds 0.5 hours but yields 2.1x additional speedup on average."
        },
        # Slide 11: Experimental Setup
        {
            "layout": 1,
            "title": "Experimental Methodology",
            "content": ["Datasets: ImageNet-1K, COCO 2017, ADE20K",
                        "Hardware: Cortex-A76, RISC-V Xuantie C910, Xilinx ZCU104 FPGA",
                        "Baselines: EfficientNet-B0/B1, MobileNetV3, FBNetV3",
                        "Metrics: Top-1 accuracy, latency (ms), power (mW), FLOPs"],
            "notes": "We use standard benchmarks for reproducibility. All experiments run 3 times with different seeds."
        },
        # Slide 12: ImageNet Results
        {
            "layout": 1,
            "title": "Results: ImageNet Classification",
            "content": ["EdgeNAS-Pro-S: 76.8% top-1, 4.2ms on Cortex-A76",
                        "EdgeNAS-Pro-M: 79.3% top-1, 7.1ms on Cortex-A76",
                        "EdgeNAS-Pro-L: 81.1% top-1, 9.8ms on Cortex-A76",
                        "Outperforms EfficientNet-B0 by 1.4% at same latency"],
            "notes": "Highlight the S variant — best in class for sub-5ms inference. Show comparison table."
        },
        # Slide 13: COCO Detection Results
        {
            "layout": 1,
            "title": "Results: Object Detection (COCO)",
            "content": ["EdgeNAS-Pro + SSDLite: 28.7 mAP, 12.3ms",
                        "EdgeNAS-Pro + YOLO head: 31.2 mAP, 15.8ms",
                        "MobileNetV3 + SSDLite baseline: 22.0 mAP, 14.1ms",
                        "6.7 mAP improvement with 13% latency reduction"],
            "notes": "Detection results demonstrate backbone quality transfers to downstream tasks."
        },
        # Slide 14: Cross-Platform Results
        {
            "layout": 1,
            "title": "Cross-Platform Performance",
            "content": ["RISC-V: 73.9% top-1, 18.4ms (3.2x faster than baseline)",
                        "FPGA: 78.1% top-1, 2.1ms (custom dataflow implementation)",
                        "Power efficiency: 340 mW (ARM), 280 mW (FPGA)",
                        "Single search produces competitive models for all 3 platforms"],
            "notes": "FPGA results are particularly impressive. Discuss custom operator library we developed."
        },
        # Slide 15: Ablation Study - Search Phases
        {
            "layout": 1,
            "title": "Ablation Study: Search Phases",
            "content": ["Macro only: 74.2% top-1 (Phase 1 alone)",
                        "Macro + Micro: 79.1% top-1 (Phase 1+2)",
                        "Full pipeline: 81.1% top-1 (Phase 1+2+3)",
                        "Each phase contributes statistically significant gains (p<0.01)"],
            "notes": "The three-phase design is validated. Removing any phase degrades results significantly."
        },
        # Slide 16: Ablation Study - Hardware Awareness
        {
            "layout": 1,
            "title": "Ablation Study: Hardware-Aware Search",
            "content": ["Without HW awareness: 78.4% top-1, 11.2ms latency",
                        "With HW awareness: 79.3% top-1, 7.1ms latency",
                        "36.6% latency reduction with 0.9% accuracy gain",
                        "Hardware lookup tables add only 12 minutes to search"],
            "notes": "Hardware awareness is essentially free in terms of search cost but critical for deployment."
        },
        # Slide 17: Search Cost Comparison
        {
            "layout": 1,
            "title": "Search Cost Analysis",
            "content": ["EdgeNAS-Pro: 4.1 GPU hours (single V100)",
                        "EfficientNet NAS: 3,800 GPU hours",
                        "Once-for-All: 40 GPU hours (amortized)",
                        "Our method: 925x cheaper than original NAS"],
            "notes": "This is our strongest selling point. Environmental impact: 47kg CO2 vs 43,560kg CO2 for full NAS."
        },
        # Slide 18: Case Study - Autonomous Drone
        {
            "layout": 1,
            "title": "Case Study: Autonomous Drone Navigation",
            "content": ["Partner: AeroVision Robotics, field deployment in Nairobi",
                        "Task: Real-time obstacle detection at 30 FPS",
                        "EdgeNAS-Pro model: 94.2% obstacle recall, 5.8ms per frame",
                        "Battery life improvement: 23% vs hand-designed model"],
            "notes": "Real-world validation is crucial. AeroVision provided flight data from 200+ missions."
        },
        # Slide 19: Case Study - Medical Imaging
        {
            "layout": 1,
            "title": "Case Study: Point-of-Care Medical Imaging",
            "content": ["Partner: Médecins Sans Frontières, deployment in rural clinics",
                        "Task: Tuberculosis screening from chest X-rays",
                        "EdgeNAS-Pro model: 96.7% sensitivity, 94.3% specificity",
                        "Runs on Raspberry Pi 4 with 180ms inference time"],
            "notes": "This application demonstrates social impact. Model is now deployed in 12 clinics across Sub-Saharan Africa."
        },
        # Slide 20: Scalability Analysis
        {
            "layout": 1,
            "title": "Scalability and Generalization",
            "content": ["Tested on 7 additional edge platforms post-hoc",
                        "Average accuracy within 1.2% of platform-specific search",
                        "Search cost scales linearly with number of target platforms",
                        "Framework supports plugin architecture for new hardware"],
            "notes": "Generalization is key for industry adoption. Discuss the plugin system briefly."
        },
        # Slide 21: Limitations
        {
            "layout": 1,
            "title": "Limitations and Threats to Validity",
            "content": ["Search space limited to CNN and lightweight attention blocks",
                        "GPU hours measured on V100 — results may vary on other GPUs",
                        "FPGA implementation requires manual operator tuning",
                        "Long-term deployment stability not yet studied"],
            "notes": "Be honest about limitations. Reviewers appreciate transparency."
        },
        # Slide 22: Future Work
        {
            "layout": 1,
            "title": "Future Directions",
            "content": ["Extend to large language model compression for edge",
                        "Incorporate dynamic inference with early exit mechanisms",
                        "Develop federated NAS for privacy-preserving search",
                        "Establish open benchmark suite for edge NAS evaluation"],
            "notes": "Future work section — mention we have a grant proposal pending for the federated NAS direction."
        },
        # Slide 23: Key Contributions
        {
            "layout": 1,
            "title": "Summary of Contributions",
            "content": ["EdgeNAS-Pro: First multi-platform NAS under 5 GPU hours",
                        "Hardware-aware search with negligible overhead",
                        "State-of-the-art accuracy-latency tradeoff on 3 platforms",
                        "Open-source framework and pretrained model zoo"],
            "notes": "Recap the four main contributions. Emphasize open-source availability."
        },
        # Slide 24: Acknowledgments
        {
            "layout": 1,
            "title": "Acknowledgments",
            "content": ["National Science Foundation — Grant #2024-1847",
                        "Horizon Europe — Project EDGE-AI-2024",
                        "AeroVision Robotics and Médecins Sans Frontières",
                        "Singapore National Supercomputing Centre for compute resources"],
            "notes": "Thank the funding bodies and collaborators. Mention student contributors."
        },
        # Slide 25: Q&A
        {
            "layout": 0,
            "title": "Questions & Discussion",
            "subtitle": "Dr. Amara Okafor — amara.okafor@ntu.edu.sg\nProf. Hiroshi Tanaka — h.tanaka@osaka-u.ac.jp\nDr. Elena Vasquez — e.vasquez@ethz.ch\n\nCode: github.com/edgenas-pro/framework",
            "notes": "Open the floor for questions. Have backup slides ready for deep-dive topics on search space design and FPGA implementation."
        },
    ]

    for i, sd in enumerate(slides_data):
        layout_idx = sd["layout"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = sd["title"]
            for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                run.font.name = "Calibri"
                if layout_idx == 0:
                    run.font.size = Pt(36)
                else:
                    run.font.size = Pt(28)
                run.font.bold = True

        # Set content
        if "subtitle" in sd and layout_idx == 0:
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = sd["subtitle"]
                for para in slide.placeholders[1].text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = "Calibri"
                        run.font.size = Pt(18)

        if "content" in sd and layout_idx == 1:
            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                add_bullet_content(tf, sd["content"])

        # Add speaker notes
        if "notes" in sd:
            slide.notes_slide.notes_text_frame.text = sd["notes"]

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
