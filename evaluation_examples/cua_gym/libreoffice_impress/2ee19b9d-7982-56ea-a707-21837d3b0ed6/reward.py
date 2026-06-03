"""
Reward Script: Configure multi-monitor presenter view in research_symposium.pptx
Task ID: impress_gf5_041
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): showPr element with presenter/speaker mode (<p:present/>)
  Component 2 (0.35): Presenter Console enabled (presOn val="1")
  Component 3 (0.25): Timer annotation in slide 1 notes (count up, 20-min warning)
  Component 4 (0.15): showAll and useTimings="0" settings present
"""

import os
import zipfile
import xml.etree.ElementTree as ET
import re

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_041'


def persist_app_state(domain: str):
    """Best-effort save of any unsaved LibreOffice edits."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def _find_presOn_val(showPr, ns):
    """Search showPr extensions for presOn element and return its val attribute."""
    for ext in showPr.findall('.//p:ext', ns):
        for child in ext:
            local_name = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if local_name == 'presOn':
                return child.get('val', '')
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Namespaces used in OOXML presentation files
    ns = {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main',
        'p15': 'http://schemas.microsoft.com/office/powerpoint/2012/main',
    }

    # Load and parse the presentation XML from the pptx ZIP
    try:
        zf = zipfile.ZipFile(file_path, 'r')
        pres_xml = zf.read('ppt/presentation.xml')
        root = ET.fromstring(pres_xml)
    except Exception as e:
        print(f"CRITICAL: Cannot load/parse pptx file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Find the showPr element (Slide Show Properties)
    showPr = root.find('.//p:showPr', ns)

    # Component 1: showPr exists with <p:present/> (speaker/full-screen mode) (0.25 points)
    try:
        if showPr is not None:
            present_el = showPr.find('p:present', ns)
            if present_el is not None:
                print(f"PASS: Component 1 — showPr with <p:present/> (speaker mode) found (0.25 pts)")
                total_score += 0.25
            else:
                # Also accept if showPr exists but has no type child (default is speaker mode)
                # Check it doesn't have <p:browse/> or <p:kiosk/> instead
                browse = showPr.find('p:browse', ns)
                kiosk = showPr.find('p:kiosk', ns)
                if browse is None and kiosk is None:
                    print(f"PASS: Component 1 — showPr exists, no browse/kiosk (defaults to speaker mode) (0.25 pts)")
                    total_score += 0.25
                else:
                    print(f"FAIL: Component 1 — showPr found but not in speaker mode (browse={browse is not None}, kiosk={kiosk is not None})")
        else:
            print(f"FAIL: Component 1 — No showPr element found in presentation.xml")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Presenter Console enabled — presOn val="1" (0.35 points)
    try:
        pres_on_val = _find_presOn_val(showPr, ns) if showPr is not None else None

        if pres_on_val == '1':
            print(f"PASS: Component 2 — Presenter Console enabled (presOn val='1') (0.35 pts)")
            total_score += 0.35
        else:
            print(f"FAIL: Component 2 — presOn val='1' not found in showPr extensions (got: {pres_on_val})")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Timer annotation in slide 1 notes — must contain count-up and 20-min warning (0.25 points)
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slide1_notes = ''
        try:
            slide1_notes = prs.slides[0].notes_slide.notes_text_frame.text
        except Exception:
            pass

        # Check for timer annotation: must mention "timer" and "20" (minutes) and count-up concept
        notes_lower = slide1_notes.lower()
        has_timer_ref = 'timer' in notes_lower
        has_20_min = bool(re.search(r'20[:.]?00|20\s*min', notes_lower))
        has_count_up = 'count up' in notes_lower or 'count-up' in notes_lower or '0:00' in slide1_notes

        if has_timer_ref and has_20_min and has_count_up:
            print(f"PASS: Component 3 — Slide 1 notes contain timer annotation with count-up and 20-min warning (0.25 pts)")
            print(f"  Notes excerpt: ...{slide1_notes[-120:]}")
            total_score += 0.25
        elif has_timer_ref and has_20_min:
            # Partial: has timer and 20 min but missing explicit count-up
            print(f"PARTIAL: Component 3 — Timer + 20-min found but count-up reference unclear (0.15 pts)")
            total_score += 0.15
        elif has_20_min:
            # Minimal: just the 20-minute reference
            print(f"PARTIAL: Component 3 — 20-minute reference found but missing timer annotation (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 3 — No timer annotation in slide 1 notes")
            print(f"  Slide 1 notes: {slide1_notes[:200]}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: showAll and useTimings="0" in showPr (0.15 points)
    try:
        if showPr is not None:
            show_all = showPr.find('p:showAll', ns)
            use_timings = showPr.get('useTimings', None)

            has_show_all = show_all is not None
            # useTimings="0" means don't auto-advance (manual control, which is normal for presenter mode)
            has_use_timings_off = use_timings == '0'

            if has_show_all and has_use_timings_off:
                print(f"PASS: Component 4 — showAll present and useTimings='0' (0.15 pts)")
                total_score += 0.15
            elif has_show_all or has_use_timings_off:
                print(f"PARTIAL: Component 4 — showAll={has_show_all}, useTimings='0'={has_use_timings_off} (0.08 pts)")
                total_score += 0.08
            else:
                print(f"FAIL: Component 4 — Neither showAll nor useTimings='0' found")
        else:
            print(f"FAIL: Component 4 — No showPr element, cannot check showAll/useTimings")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    try:
        zf.close()
    except:
        pass

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state('libreoffice_impress')

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
