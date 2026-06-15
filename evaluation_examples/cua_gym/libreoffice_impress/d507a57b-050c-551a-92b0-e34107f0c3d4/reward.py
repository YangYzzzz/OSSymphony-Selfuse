"""
Reward Script: Tokyo Travel Guide Presentation
Task ID: impress_wf_047
Domain: libreoffice_impress
Scoring:
  C1 (0.15) - 12 slides
  C2 (0.10) - Slide 1 title 'Tokyo Travel Guide' with red #C62828 accent
  C3 (0.10) - Slide 2 colored circles (TOC with icon circles)
  C4 (0.10) - Slides 3-4 attractions with image placeholders
  C5 (0.10) - Slides 5-6 food guide tables
  C6 (0.15) - Slide 7 transit diagram (circles + lines)
  C7 (0.10) - Slide 8 accommodation comparison table
  C8 (0.05) - Slide 9 pie chart
  C9 (0.10) - Slide 10 green do's and red don'ts
  C10 (0.05) - Slide 11 itinerary table
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_047'


def persist_app_state(domain):
    """Attempt to save any unsaved GUI state."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Presentation has 12 slides (0.15 pts)
    try:
        if num_slides == 12:
            print(f"PASS: Component 1 — 12 slides found (0.15 pts)")
            total_score += 0.15
        elif num_slides >= 10:
            partial = 0.08
            print(f"PARTIAL: Component 1 — {num_slides} slides (expected 12), awarding {partial} pts")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — found {num_slides} slides, expected 12")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    if num_slides < 2:
        # Can't check further components meaningfully
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: Slide 1 title contains 'Tokyo Travel Guide' and has red C62828 accent (0.10 pts)
    try:
        slide1 = prs.slides[0]
        has_title = False
        has_red_accent = False
        for shape in slide1.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip().lower()
                if 'tokyo travel guide' in txt:
                    has_title = True
            # Check for C62828 fill color on shapes
            try:
                fill = shape.fill
                if fill.type == 1:  # solid fill
                    rgb_str = str(fill.fore_color.rgb).upper()
                    if rgb_str == 'C62828':
                        has_red_accent = True
            except:
                pass
        if has_title and has_red_accent:
            print(f"PASS: Component 2 — Title 'Tokyo Travel Guide' with red #C62828 accent (0.10 pts)")
            total_score += 0.10
        elif has_title:
            print(f"PARTIAL: Component 2 — Title found but no red #C62828 accent shape (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 2 — Title 'Tokyo Travel Guide' not found on slide 1")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 2 has colored circles (TOC icons) — at least 4 auto shapes with solid fills (0.10 pts)
    try:
        if num_slides >= 2:
            slide2 = prs.slides[1]
            colored_circles = 0
            for shape in slide2.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    try:
                        fill = shape.fill
                        if fill.type == 1:  # solid fill
                            colored_circles += 1
                    except:
                        pass
            if colored_circles >= 4:
                print(f"PASS: Component 3 — Slide 2 has {colored_circles} colored circle shapes (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 3 — Slide 2 has only {colored_circles} colored auto shapes, expected >= 4")
        else:
            print(f"FAIL: Component 3 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slides 3-4 have image placeholders and attraction descriptions (0.10 pts)
    try:
        if num_slides >= 4:
            attractions_ok = 0
            for idx in [2, 3]:  # slides 3 and 4
                slide = prs.slides[idx]
                has_placeholder = False
                has_description = False
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        txt = shape.text_frame.text.strip()
                        if '[photo]' in txt.lower() or '[image]' in txt.lower():
                            has_placeholder = True
                        if len(txt) > 40:  # description text
                            has_description = True
                if has_placeholder and has_description:
                    attractions_ok += 1
            if attractions_ok == 2:
                print(f"PASS: Component 4 — Slides 3-4 both have image placeholders and descriptions (0.10 pts)")
                total_score += 0.10
            elif attractions_ok == 1:
                print(f"PARTIAL: Component 4 — Only 1 of 2 attraction slides fully correct (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 4 — Attraction slides missing placeholders or descriptions")
        else:
            print(f"FAIL: Component 4 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slides 5-6 have food guide tables (0.10 pts)
    try:
        if num_slides >= 6:
            table_count = 0
            for idx in [4, 5]:  # slides 5 and 6
                slide = prs.slides[idx]
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        t = shape.table
                        if len(t.rows) >= 2 and len(t.columns) >= 2:
                            table_count += 1
            if table_count >= 2:
                print(f"PASS: Component 5 — Slides 5-6 each have a table (0.10 pts)")
                total_score += 0.10
            elif table_count == 1:
                print(f"PARTIAL: Component 5 — Only 1 table found on slides 5-6 (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 5 — No tables found on slides 5-6")
        else:
            print(f"FAIL: Component 5 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Slide 7 has transit diagram — circles connected by lines (0.15 pts)
    try:
        if num_slides >= 7:
            slide7 = prs.slides[6]
            line_count = 0
            circle_count = 0
            for shape in slide7.shapes:
                st = str(shape.shape_type)
                if 'LINE' in st:
                    line_count += 1
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    circle_count += 1
            if line_count >= 5 and circle_count >= 5:
                print(f"PASS: Component 6 — Slide 7 transit diagram: {circle_count} circles, {line_count} lines (0.15 pts)")
                total_score += 0.15
            elif line_count >= 3 and circle_count >= 3:
                print(f"PARTIAL: Component 6 — Slide 7 partial diagram: {circle_count} circles, {line_count} lines (0.08 pts)")
                total_score += 0.08
            else:
                print(f"FAIL: Component 6 — Slide 7 missing transit diagram (circles={circle_count}, lines={line_count})")
        else:
            print(f"FAIL: Component 6 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Slide 8 has accommodation comparison table (0.10 pts)
    try:
        if num_slides >= 8:
            slide8 = prs.slides[7]
            has_table = False
            for shape in slide8.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    t = shape.table
                    if len(t.rows) >= 3 and len(t.columns) >= 3:
                        has_table = True
            if has_table:
                print(f"PASS: Component 7 — Slide 8 has accommodation comparison table (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 7 — Slide 8 missing comparison table")
        else:
            print(f"FAIL: Component 7 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    # Component 8: Slide 9 has a pie chart (0.05 pts)
    try:
        if num_slides >= 9:
            slide9 = prs.slides[8]
            has_chart = False
            for shape in slide9.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    # Check if it's a pie chart (type 5)
                    try:
                        if shape.chart.chart_type == 5:  # PIE
                            has_chart = True
                        else:
                            has_chart = True  # any chart is partial credit
                    except:
                        has_chart = True
            if has_chart:
                print(f"PASS: Component 8 — Slide 9 has pie chart (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 8 — Slide 9 missing chart")
        else:
            print(f"FAIL: Component 8 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 8 — {e}")

    # Component 9: Slide 10 has green do's and red don'ts text (0.10 pts)
    try:
        if num_slides >= 10:
            slide10 = prs.slides[9]
            has_green_text = False
            has_red_text = False
            for shape in slide10.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                if run.font.color.type is not None:
                                    rgb_str = str(run.font.color.rgb).upper()
                                    # Green family check
                                    r = int(rgb_str[0:2], 16)
                                    g = int(rgb_str[2:4], 16)
                                    b = int(rgb_str[4:6], 16)
                                    if g > r and g > b and g >= 100:
                                        has_green_text = True
                                    # Red family check (C62828, or any red-dominant)
                                    if r > g and r > b and r >= 150:
                                        txt = run.text.strip().lower()
                                        # Make sure it's in the don'ts section, not the title
                                        if "don" in txt or "don't" in txt or "tip" in txt or "blow" in txt or "stick" in txt or "pour" in txt or "eat" in txt or "talk" in txt:
                                            has_red_text = True
                                        elif "don" in txt.lower():
                                            has_red_text = True
                            except:
                                pass
            # Simpler red check: any text with red color on this slide that isn't the title
            if not has_red_text:
                for shape in slide10.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                try:
                                    if run.font.color.type is not None:
                                        rgb_str = str(run.font.color.rgb).upper()
                                        if rgb_str in ('C62828', 'FF0000', 'CC0000', 'B71C1C', 'D32F2F', 'E53935'):
                                            txt = run.text.strip()
                                            if "don" in txt.lower() or len(txt) > 10:
                                                has_red_text = True
                                except:
                                    pass

            if has_green_text and has_red_text:
                print(f"PASS: Component 9 — Slide 10 has green do's and red don'ts (0.10 pts)")
                total_score += 0.10
            elif has_green_text or has_red_text:
                print(f"PARTIAL: Component 9 — green={has_green_text}, red={has_red_text} (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 9 — Slide 10 missing colored do's/don'ts")
        else:
            print(f"FAIL: Component 9 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 9 — {e}")

    # Component 10: Slide 11 has itinerary table (0.05 pts)
    try:
        if num_slides >= 11:
            slide11 = prs.slides[10]
            has_table = False
            for shape in slide11.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    t = shape.table
                    if len(t.rows) >= 2:
                        has_table = True
            if has_table:
                print(f"PASS: Component 10 — Slide 11 has itinerary table (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 10 — Slide 11 missing itinerary table")
        else:
            print(f"FAIL: Component 10 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 10 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
