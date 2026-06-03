"""
Reward Script: Configure master slide with two-column content layout
Task ID: impress_gf3_036
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Master slide has exactly two body/content placeholders (non-title, non-date/footer/slidenum)
  Component 2 (0.25): Both content placeholders are ~11.50cm wide
  Component 3 (0.15): Both content placeholders are vertically aligned at same top position
  Component 4 (0.15): Gap between the two placeholders is ~0.50cm
  Component 5 (0.15): Title placeholder is unchanged (same position/size as original)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_036'

# EMU constants
EMU_PER_CM = 360000  # 1 cm = 360000 EMU


def emu_to_cm(emu):
    """Convert EMU to centimeters."""
    return emu / EMU_PER_CM


def is_approx_equal(val1, val2, tolerance_cm=0.15):
    """Check if two EMU values are approximately equal within tolerance in cm."""
    return abs(emu_to_cm(val1) - emu_to_cm(val2)) <= tolerance_cm


def is_approx_cm(emu_val, expected_cm, tolerance_cm=0.15):
    """Check if an EMU value is approximately equal to expected cm value."""
    return abs(emu_to_cm(emu_val) - expected_cm) <= tolerance_cm


def persist_app_state(domain):
    """Save any unsaved GUI state before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(1.0)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Get the first slide master
    if len(prs.slide_masters) == 0:
        print("CRITICAL: No slide masters found")
        print("REWARD: 0.0")
        return 0.0

    master = prs.slide_masters[0]

    # Categorize placeholders on the master slide
    title_phs = []
    content_phs = []  # body/content type placeholders (not title, date, footer, slidenum)
    other_phs = []

    SKIP_TYPES = {13, 15, 16}  # SLIDE_NUMBER, FOOTER, DATE

    for ph in master.placeholders:
        ph_type = ph.placeholder_format.type
        ph_type_val = int(ph_type) if ph_type is not None else -1

        if ph_type_val == 1:  # TITLE
            title_phs.append(ph)
        elif ph_type_val in SKIP_TYPES:
            other_phs.append(ph)
        else:
            # BODY (2), OBJECT (7), or other content types
            content_phs.append(ph)

    print(f"INFO: Master slide has {len(title_phs)} title placeholder(s), "
          f"{len(content_phs)} content placeholder(s), {len(other_phs)} other placeholder(s)")

    # Component 1: Master slide has exactly two body/content placeholders (0.30 points)
    try:
        if len(content_phs) == 2:
            print(f"PASS: Component 1 — Master has exactly 2 content placeholders (0.30 pts)")
            total_score += 0.30
        elif len(content_phs) > 2:
            print(f"FAIL: Component 1 — Master has {len(content_phs)} content placeholders, expected 2")
        else:
            print(f"FAIL: Component 1 — Master has {len(content_phs)} content placeholders, expected 2")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Early exit if we don't have exactly 2 content placeholders
    if len(content_phs) != 2:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Sort content placeholders by left position (left one first)
    content_phs.sort(key=lambda ph: ph.left)
    left_ph = content_phs[0]
    right_ph = content_phs[1]

    print(f"INFO: Left placeholder — left={emu_to_cm(left_ph.left):.2f}cm, "
          f"width={emu_to_cm(left_ph.width):.2f}cm, top={emu_to_cm(left_ph.top):.2f}cm, "
          f"height={emu_to_cm(left_ph.height):.2f}cm")
    print(f"INFO: Right placeholder — left={emu_to_cm(right_ph.left):.2f}cm, "
          f"width={emu_to_cm(right_ph.width):.2f}cm, top={emu_to_cm(right_ph.top):.2f}cm, "
          f"height={emu_to_cm(right_ph.height):.2f}cm")

    # Component 2: Both content placeholders are ~11.50cm wide (0.25 points)
    try:
        left_width_cm = emu_to_cm(left_ph.width)
        right_width_cm = emu_to_cm(right_ph.width)
        left_ok = is_approx_cm(left_ph.width, 11.50, tolerance_cm=0.30)
        right_ok = is_approx_cm(right_ph.width, 11.50, tolerance_cm=0.30)

        if left_ok and right_ok:
            print(f"PASS: Component 2 — Both placeholders ~11.50cm wide "
                  f"(left={left_width_cm:.2f}cm, right={right_width_cm:.2f}cm) (0.25 pts)")
            total_score += 0.25
        elif left_ok or right_ok:
            # Partial: one is correct
            print(f"PARTIAL: Component 2 — Only one placeholder ~11.50cm wide "
                  f"(left={left_width_cm:.2f}cm, right={right_width_cm:.2f}cm) (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 — Neither placeholder is ~11.50cm wide "
                  f"(left={left_width_cm:.2f}cm, right={right_width_cm:.2f}cm)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Both placeholders vertically aligned at same top position (0.15 points)
    try:
        tops_match = is_approx_equal(left_ph.top, right_ph.top, tolerance_cm=0.20)
        heights_match = is_approx_equal(left_ph.height, right_ph.height, tolerance_cm=0.30)

        if tops_match and heights_match:
            print(f"PASS: Component 3 — Both placeholders vertically aligned "
                  f"(top_left={emu_to_cm(left_ph.top):.2f}cm, top_right={emu_to_cm(right_ph.top):.2f}cm, "
                  f"h_left={emu_to_cm(left_ph.height):.2f}cm, h_right={emu_to_cm(right_ph.height):.2f}cm) (0.15 pts)")
            total_score += 0.15
        elif tops_match:
            print(f"PARTIAL: Component 3 — Tops aligned but heights differ "
                  f"(h_left={emu_to_cm(left_ph.height):.2f}cm, h_right={emu_to_cm(right_ph.height):.2f}cm) (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 3 — Placeholders not vertically aligned "
                  f"(top_left={emu_to_cm(left_ph.top):.2f}cm, top_right={emu_to_cm(right_ph.top):.2f}cm)")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Gap between the two placeholders is ~0.50cm (0.15 points)
    try:
        gap_emu = right_ph.left - (left_ph.left + left_ph.width)
        gap_cm = emu_to_cm(gap_emu)

        if abs(gap_cm - 0.50) <= 0.20:
            print(f"PASS: Component 4 — Gap between placeholders is {gap_cm:.2f}cm (~0.50cm) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 — Gap between placeholders is {gap_cm:.2f}cm, expected ~0.50cm")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Title placeholder is unchanged (0.15 points)
    # Initial title: left=457200 (1.27cm), top=274638 (0.76cm), width=8229600 (22.86cm), height=1143000 (3.17cm)
    try:
        if len(title_phs) == 1:
            title = title_phs[0]
            title_left_ok = is_approx_cm(title.left, 1.27, tolerance_cm=0.20)
            title_top_ok = is_approx_cm(title.top, 0.76, tolerance_cm=0.20)
            title_width_ok = is_approx_cm(title.width, 22.86, tolerance_cm=0.30)
            title_height_ok = is_approx_cm(title.height, 3.17, tolerance_cm=0.30)

            if title_left_ok and title_top_ok and title_width_ok and title_height_ok:
                print(f"PASS: Component 5 — Title placeholder unchanged "
                      f"(left={emu_to_cm(title.left):.2f}cm, top={emu_to_cm(title.top):.2f}cm, "
                      f"w={emu_to_cm(title.width):.2f}cm, h={emu_to_cm(title.height):.2f}cm) (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 5 — Title placeholder changed "
                      f"(left={emu_to_cm(title.left):.2f}cm, top={emu_to_cm(title.top):.2f}cm, "
                      f"w={emu_to_cm(title.width):.2f}cm, h={emu_to_cm(title.height):.2f}cm)")
        else:
            print(f"FAIL: Component 5 — Expected 1 title placeholder, found {len(title_phs)}")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist app state before verification (LibreOffice may have unsaved changes)
persist_app_state('libreoffice_impress')

# Test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
