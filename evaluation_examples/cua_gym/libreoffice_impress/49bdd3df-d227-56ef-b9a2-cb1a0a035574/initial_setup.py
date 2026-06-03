"""
Initial Setup: Configure master slide two-column layout task
Task ID: impress_gf3_036
Domain: libreoffice_impress

Creates a 10-slide presentation with a master slide that has a title placeholder
and a single full-width content placeholder. Slides contain realistic business content.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_036'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard 16:9 slide dimensions
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(19.05)

    # We'll use layout index 1 (Title + Content) which has a title and a single
    # full-width content placeholder on the master - this matches the task requirement.
    # The default template already has a slide master with title + single content placeholder.

    layout_title_content = prs.slide_layouts[1]  # Title + Content layout

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q3 2025 Strategic Review"
    slide1.placeholders[1].text = "Global Operations Division\nPrepared by Sarah Chen, VP Strategy"

    # Slide 2: Agenda
    slide2 = prs.slides.add_slide(layout_title_content)
    slide2.shapes.title.text = "Meeting Agenda"
    tf = slide2.placeholders[1].text_frame
    tf.text = "1. Financial Performance Summary"
    for item in [
        "2. Regional Market Analysis",
        "3. Product Pipeline Updates",
        "4. Customer Acquisition Metrics",
        "5. Technology Infrastructure Roadmap",
        "6. Talent & Workforce Planning",
        "7. Risk Assessment & Mitigation",
        "8. Q4 Strategic Priorities",
    ]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

    # Slide 3: Financial Summary
    slide3 = prs.slides.add_slide(layout_title_content)
    slide3.shapes.title.text = "Financial Performance"
    tf = slide3.placeholders[1].text_frame
    tf.text = "Revenue: $142.8M (+12.3% YoY)"
    items = [
        "Operating Margin: 23.7% (up from 21.2%)",
        "EBITDA: $33.8M exceeding target by $2.1M",
        "Free Cash Flow: $28.4M, strongest quarter since 2023",
        "SG&A Expenses: $18.2M (12.7% of revenue)",
        "R&D Investment: $14.6M (10.2% of revenue)",
    ]
    for item in items:
        p = tf.add_paragraph()
        p.text = item

    # Slide 4: Regional Analysis
    slide4 = prs.slides.add_slide(layout_title_content)
    slide4.shapes.title.text = "Regional Market Analysis"
    tf = slide4.placeholders[1].text_frame
    tf.text = "North America: $68.3M (47.8% share, +8.1%)"
    for item in [
        "Europe: $38.7M (27.1% share, +15.4%)",
        "Asia-Pacific: $24.9M (17.4% share, +22.6%)",
        "Latin America: $7.2M (5.0% share, +31.2%)",
        "Middle East & Africa: $3.7M (2.6% share, +18.9%)",
    ]:
        p = tf.add_paragraph()
        p.text = item

    # Slide 5: Product Pipeline
    slide5 = prs.slides.add_slide(layout_title_content)
    slide5.shapes.title.text = "Product Pipeline Updates"
    tf = slide5.placeholders[1].text_frame
    tf.text = "Aurora Platform v3.2 - Released Aug 15, adoption at 73%"
    for item in [
        "Nexus Analytics Suite - Beta testing, 142 enterprise clients",
        "CloudBridge Integration - On track for Oct launch",
        "Mobile SDK 2.0 - Delayed 3 weeks, new ETA: Nov 8",
        "AI-Powered Insights Module - Prototype approved by steering committee",
    ]:
        p = tf.add_paragraph()
        p.text = item

    # Slide 6: Customer Metrics
    slide6 = prs.slides.add_slide(layout_title_content)
    slide6.shapes.title.text = "Customer Acquisition & Retention"
    tf = slide6.placeholders[1].text_frame
    tf.text = "New Enterprise Clients: 47 (+28% vs Q2)"
    for item in [
        "Net Revenue Retention: 118% (target: 115%)",
        "Customer Acquisition Cost: $12,400 (down 15%)",
        "Lifetime Value: $287,000 (up 9%)",
        "Churn Rate: 3.2% (industry avg: 5.8%)",
        "NPS Score: 72 (up from 68 in Q2)",
    ]:
        p = tf.add_paragraph()
        p.text = item

    # Slide 7: Technology
    slide7 = prs.slides.add_slide(layout_title_content)
    slide7.shapes.title.text = "Technology Infrastructure"
    tf = slide7.placeholders[1].text_frame
    tf.text = "Cloud Migration: 87% complete (target: 95% by EOY)"
    for item in [
        "Platform Uptime: 99.97% (SLA: 99.95%)",
        "API Response Time: 42ms avg (down from 67ms)",
        "Security Incidents: 0 critical, 3 minor (all resolved)",
        "Tech Debt Reduction: 23% of backlog cleared",
    ]:
        p = tf.add_paragraph()
        p.text = item

    # Slide 8: Talent
    slide8 = prs.slides.add_slide(layout_title_content)
    slide8.shapes.title.text = "Talent & Workforce Planning"
    tf = slide8.placeholders[1].text_frame
    tf.text = "Headcount: 1,247 (+89 net new hires)"
    for item in [
        "Engineering Team: 412 (33% of workforce)",
        "Voluntary Turnover: 8.4% annualized (industry: 13%)",
        "Offer Acceptance Rate: 91%",
        "Diversity Hiring: 48% of new hires from underrepresented groups",
        "Employee Satisfaction Score: 4.3/5.0",
    ]:
        p = tf.add_paragraph()
        p.text = item

    # Slide 9: Risk Assessment
    slide9 = prs.slides.add_slide(layout_title_content)
    slide9.shapes.title.text = "Risk Assessment"
    tf = slide9.placeholders[1].text_frame
    tf.text = "HIGH: Supply chain disruption in APAC region"
    for item in [
        "MEDIUM: Regulatory changes in EU data privacy framework",
        "MEDIUM: Competitor pricing pressure in SMB segment",
        "LOW: Currency fluctuation impact on LATAM revenue",
        "MITIGATED: Vendor concentration risk (diversified Q2)",
    ]:
        p = tf.add_paragraph()
        p.text = item

    # Slide 10: Q4 Priorities
    slide10 = prs.slides.add_slide(layout_title_content)
    slide10.shapes.title.text = "Q4 Strategic Priorities"
    tf = slide10.placeholders[1].text_frame
    tf.text = "1. Complete CloudBridge launch and onboard 30 beta clients"
    for item in [
        "2. Achieve 95% cloud migration milestone",
        "3. Expand APAC sales team by 15 representatives",
        "4. Launch customer success program for top-tier accounts",
        "5. Finalize 2026 budget and strategic plan",
        "6. Close 3 strategic partnership agreements",
    ]:
        p = tf.add_paragraph()
        p.text = item

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
