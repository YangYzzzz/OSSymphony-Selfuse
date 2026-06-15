"""
Initial Setup: Create a 15-slide conference presentation with no sections defined.
Task ID: impress_ndo_071
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_071'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # Conference presentation slide content
    slide_contents = [
        # Slides 1-3: Introduction-style
        {"layout": 0, "title": "Annual Technology Conference 2025",
         "subtitle": "Innovations Shaping Tomorrow's Enterprise"},
        {"layout": 1, "title": "Conference Agenda",
         "body": "Morning Keynote: AI in Enterprise\nPanel: Cloud-Native Architecture\nWorkshop: DevOps Best Practices\nNetworking Lunch\nAfternoon Sessions: Deep Dives\nClosing Remarks"},
        {"layout": 1, "title": "Speaker Overview",
         "body": "Dr. Elena Vasquez - Chief AI Researcher, NovaTech\nRaj Patel - VP Engineering, CloudScale Inc.\nSarah Kim - Director of Platform Engineering, ByteForge\nMarcus Chen - CTO, DataStream Analytics\nAisha Okafor - Head of DevOps, QuantumLeap"},

        # Slides 4-10: Main Content
        {"layout": 1, "title": "The State of AI in 2025",
         "body": "Large Language Models driving automation\nMultimodal AI transforming user interfaces\nEdge AI reducing latency by 40%\nResponsible AI frameworks gaining adoption\nAI spending projected at $200B globally"},
        {"layout": 1, "title": "Cloud-Native Architecture Trends",
         "body": "Kubernetes adoption at 78% among enterprises\nServerless computing growing 25% year-over-year\nMulti-cloud strategies becoming standard\nService mesh adoption tripling since 2023\nFinOps practices reducing cloud waste by 30%"},
        {"layout": 1, "title": "DevOps Maturity Model",
         "body": "Level 1: Initial - Ad hoc processes\nLevel 2: Managed - Basic CI/CD pipelines\nLevel 3: Defined - Infrastructure as Code\nLevel 4: Measured - Full observability\nLevel 5: Optimized - Self-healing systems"},
        {"layout": 1, "title": "Security Best Practices",
         "body": "Zero Trust Architecture implementation\nShift-left security in CI/CD pipelines\nSecret management with HashiCorp Vault\nContainer image scanning with Trivy\nSupply chain security via SBOM generation"},
        {"layout": 1, "title": "Platform Engineering Workshop",
         "body": "Internal Developer Portals with Backstage\nSelf-service infrastructure provisioning\nGolden paths for microservice deployment\nDeveloper experience metrics and KPIs\nReducing cognitive load through abstraction"},
        {"layout": 1, "title": "Data Engineering at Scale",
         "body": "Real-time streaming with Apache Kafka\nData lakehouse architecture patterns\nMLOps pipeline orchestration\nData quality monitoring frameworks\nGovernance and compliance automation"},
        {"layout": 1, "title": "Case Study: Enterprise Migration",
         "body": "Client: Global Financial Services Corp\nChallenge: Legacy monolith to microservices\nTimeline: 18 months, 200+ engineers\nResult: 60% reduction in deployment time\nROI: $12M annual infrastructure savings"},

        # Slides 11-15: Conclusion
        {"layout": 1, "title": "Key Takeaways",
         "body": "AI is no longer optional - integrate or fall behind\nCloud-native is the foundation for scale\nPlatform engineering empowers developers\nSecurity must be embedded, not bolted on\nData is the new competitive advantage"},
        {"layout": 1, "title": "2025 Predictions",
         "body": "50% of enterprises will adopt AI agents\nMulti-cloud will be default architecture\nPlatform teams will outnumber ops teams\nAI-assisted coding reaches 80% adoption\nGreen computing becomes a board priority"},
        {"layout": 1, "title": "Resources and Next Steps",
         "body": "Conference recordings: techconf2025.io/recordings\nSlack community: techconf2025.slack.com\nGitHub repos: github.com/techconf2025\nNewsletter signup: techconf2025.io/newsletter\nNext conference: October 15-17, 2025 - Berlin"},
        {"layout": 1, "title": "Q&A Session",
         "body": "Submit questions via the conference app\nLive polling available at sli.do/techconf25\nBreakout rooms available for deeper discussions\nSpeaker office hours: 4:00 PM - 5:30 PM\nFeedback forms at each exit"},
        {"layout": 0, "title": "Thank You!",
         "subtitle": "See you at TechConf 2025 Berlin\nwww.techconf2025.io"},
    ]

    for i, content in enumerate(slide_contents):
        layout_idx = content["layout"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = content["title"]

        # Set body or subtitle
        if layout_idx == 0 and "subtitle" in content:
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = content["subtitle"]
        elif "body" in content:
            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                lines = content["body"].split("\n")
                for j, line in enumerate(lines):
                    if j == 0:
                        tf.paragraphs[0].text = line
                    else:
                        p = tf.add_paragraph()
                        p.text = line

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
