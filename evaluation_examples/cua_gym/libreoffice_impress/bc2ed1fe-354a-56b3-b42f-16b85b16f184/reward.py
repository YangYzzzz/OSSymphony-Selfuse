"""
Reward Script: Set different alignments across slides 2-6 in a workshop presentation
Task ID: osworld_impress_per_slide_alignment_010
Domain: libreoffice_impress

Task: Apply the following body text alignments to TextBox 3 on slides 2-6:
  - Slide 2: CENTER
  - Slide 3: RIGHT
  - Slide 4: CENTER
  - Slide 5: LEFT  (already left in initial, so no observable change)
  - Slide 6: RIGHT

Scoring Rubric (only task-introduced changes are scored):
  Component 1: Slide 2 body TextBox 3 paragraphs are CENTER aligned  (0.25 pts)
  Component 2: Slide 3 body TextBox 3 paragraphs are RIGHT aligned   (0.25 pts)
  Component 3: Slide 4 body TextBox 3 paragraphs are CENTER aligned  (0.25 pts)
  Component 4: Slide 6 body TextBox 3 paragraphs are RIGHT aligned   (0.25 pts)
  Total: 1.0

Note: Slide 5 body text is LEFT aligned in both initial and golden states (no observable change),
so it is NOT scored as a separate component to avoid awarding points for pre-existing conditions.
The horizontal re-centering of textboxes is also not scored separately since all TextBox 3 shapes
were already horizontally centered (center_offset=0.0) in the initial state.
"""

import os
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_per_slide_alignment_010'


def get_body_textbox(slide):
    """Return the 'TextBox 3' shape (the body text box) from a slide, or None if not found."""
    for shape in slide.shapes:
        if shape.name == 'TextBox 3' and shape.has_text_frame:
            return shape
    return None


def check_all_paragraphs_alignment(shape, expected_alignment):
    """
    Check that all non-empty paragraphs in a text frame have the expected alignment.
    Returns (bool, int, str) -> (all_match, para_count, detail_msg)
    """
    if not shape.has_text_frame:
        return False, 0, "Shape has no text frame"

    tf = shape.text_frame
    mismatches = []
    count = 0

    for k, para in enumerate(tf.paragraphs):
        # Skip empty paragraphs
        if not para.text.strip():
            continue
        count += 1
        actual = para.alignment
        # Normalize: None means LEFT in python-pptx
        normalized_actual = actual if actual is not None else PP_ALIGN.LEFT

        if normalized_actual != expected_alignment:
            mismatches.append(
                f"Para {k}: expected={expected_alignment}, actual={actual}"
            )

    if count == 0:
        return False, 0, "No non-empty paragraphs found"

    if mismatches:
        return False, count, f"{len(mismatches)}/{count} paragraphs have wrong alignment: {mismatches[:2]}"

    return True, count, f"All {count} paragraphs correctly aligned"


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Validate basic structure: must have at least 7 slides
    if len(prs.slides) < 7:
        print(f"CRITICAL: Expected at least 7 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # --- Component 1: Slide 2 body TextBox 3 is CENTER aligned (0.25 pts) ---
    # Initial state: LEFT -> Task requires: CENTER
    # This should FAIL on initial_env and PASS on golden_env
    try:
        slide2 = prs.slides[1]  # 0-indexed
        body2 = get_body_textbox(slide2)
        if body2 is None:
            print("FAIL: Component 1 — TextBox 3 not found on slide 2")
        else:
            ok, count, detail = check_all_paragraphs_alignment(body2, PP_ALIGN.CENTER)
            if ok:
                print(f"PASS: Component 1 — Slide 2 body text is CENTER aligned ({count} paragraphs) (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 — Slide 2 body text not CENTER aligned: {detail}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # --- Component 2: Slide 3 body TextBox 3 is RIGHT aligned (0.25 pts) ---
    # Initial state: LEFT -> Task requires: RIGHT
    # This should FAIL on initial_env and PASS on golden_env
    try:
        slide3 = prs.slides[2]  # 0-indexed
        body3 = get_body_textbox(slide3)
        if body3 is None:
            print("FAIL: Component 2 — TextBox 3 not found on slide 3")
        else:
            ok, count, detail = check_all_paragraphs_alignment(body3, PP_ALIGN.RIGHT)
            if ok:
                print(f"PASS: Component 2 — Slide 3 body text is RIGHT aligned ({count} paragraphs) (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 2 — Slide 3 body text not RIGHT aligned: {detail}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # --- Component 3: Slide 4 body TextBox 3 is CENTER aligned (0.25 pts) ---
    # Initial state: LEFT -> Task requires: CENTER
    # This should FAIL on initial_env and PASS on golden_env
    try:
        slide4 = prs.slides[3]  # 0-indexed
        body4 = get_body_textbox(slide4)
        if body4 is None:
            print("FAIL: Component 3 — TextBox 3 not found on slide 4")
        else:
            ok, count, detail = check_all_paragraphs_alignment(body4, PP_ALIGN.CENTER)
            if ok:
                print(f"PASS: Component 3 — Slide 4 body text is CENTER aligned ({count} paragraphs) (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 3 — Slide 4 body text not CENTER aligned: {detail}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # --- Component 4: Slide 6 body TextBox 3 is RIGHT aligned (0.25 pts) ---
    # Initial state: LEFT -> Task requires: RIGHT
    # This should FAIL on initial_env and PASS on golden_env
    try:
        slide6 = prs.slides[5]  # 0-indexed
        body6 = get_body_textbox(slide6)
        if body6 is None:
            print("FAIL: Component 4 — TextBox 3 not found on slide 6")
        else:
            ok, count, detail = check_all_paragraphs_alignment(body6, PP_ALIGN.RIGHT)
            if ok:
                print(f"PASS: Component 4 — Slide 6 body text is RIGHT aligned ({count} paragraphs) (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 4 — Slide 6 body text not RIGHT aligned: {detail}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
