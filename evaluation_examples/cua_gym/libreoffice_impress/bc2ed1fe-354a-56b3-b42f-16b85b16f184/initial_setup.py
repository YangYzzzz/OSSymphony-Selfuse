"""
Initial Setup: 7-slide workshop deck with left-aligned body text on slides 2-6
Task ID: osworld_impress_per_slide_alignment_010
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_per_slide_alignment_010'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_width = prs.slide_width

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Advanced Data Analytics Workshop"
    slide1.placeholders[1].text = "Transforming Raw Data into Actionable Insights\nMarch 2025"

    # --- Slide 2: Introduction ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank

    # Title textbox
    title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.9))
    tf2t = title2.text_frame
    p2t = tf2t.paragraphs[0]
    p2t.text = "Workshop Overview"
    p2t.alignment = PP_ALIGN.LEFT
    run2t = p2t.runs[0]
    run2t.font.size = Pt(32)
    run2t.font.bold = True
    run2t.font.color.rgb = RGBColor(0x1F, 0x48, 0x7E)

    # Body textbox - LEFT aligned (pre-task state, will be changed to CENTER)
    body2_width = Inches(7)
    body2_left = (slide_width - body2_width) // 2
    body2 = slide2.shapes.add_textbox(body2_left, Inches(1.4), body2_width, Inches(4.5))
    tf2 = body2.text_frame
    tf2.word_wrap = True

    lines2 = [
        "This workshop covers modern data analytics techniques and tools.",
        "Participants will learn statistical analysis fundamentals.",
        "Hands-on exercises with real-world datasets are included.",
        "Group projects encourage collaborative problem solving.",
        "Certificate of completion provided upon finishing all modules.",
    ]
    for i, line in enumerate(lines2):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 3: Module 1 ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])

    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.9))
    tf3t = title3.text_frame
    p3t = tf3t.paragraphs[0]
    p3t.text = "Module 1: Data Collection & Cleaning"
    p3t.alignment = PP_ALIGN.LEFT
    run3t = p3t.runs[0]
    run3t.font.size = Pt(32)
    run3t.font.bold = True
    run3t.font.color.rgb = RGBColor(0x1F, 0x48, 0x7E)

    # Body textbox - LEFT aligned (will be changed to RIGHT)
    body3_width = Inches(7)
    body3_left = (slide_width - body3_width) // 2
    body3 = slide3.shapes.add_textbox(body3_left, Inches(1.4), body3_width, Inches(4.5))
    tf3 = body3.text_frame
    tf3.word_wrap = True

    lines3 = [
        "Data collection methods: surveys, APIs, web scraping.",
        "Identifying and handling missing values in datasets.",
        "Removing duplicates and correcting inconsistencies.",
        "Data type validation and format standardization.",
        "Creating data quality reports for stakeholder review.",
    ]
    for i, line in enumerate(lines3):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 4: Module 2 ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])

    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.9))
    tf4t = title4.text_frame
    p4t = tf4t.paragraphs[0]
    p4t.text = "Module 2: Exploratory Data Analysis"
    p4t.alignment = PP_ALIGN.LEFT
    run4t = p4t.runs[0]
    run4t.font.size = Pt(32)
    run4t.font.bold = True
    run4t.font.color.rgb = RGBColor(0x1F, 0x48, 0x7E)

    # Body textbox - LEFT aligned (will be changed to CENTER)
    body4_width = Inches(7)
    body4_left = (slide_width - body4_width) // 2
    body4 = slide4.shapes.add_textbox(body4_left, Inches(1.4), body4_width, Inches(4.5))
    tf4 = body4.text_frame
    tf4.word_wrap = True

    lines4 = [
        "Descriptive statistics: mean, median, mode, standard deviation.",
        "Distribution analysis and histogram interpretation.",
        "Correlation matrices and scatter plot analysis.",
        "Outlier detection using IQR and Z-score methods.",
        "Visualization best practices with matplotlib and seaborn.",
    ]
    for i, line in enumerate(lines4):
        if i == 0:
            p = tf4.paragraphs[0]
        else:
            p = tf4.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 5: Module 3 ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])

    title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.9))
    tf5t = title5.text_frame
    p5t = tf5t.paragraphs[0]
    p5t.text = "Module 3: Statistical Modeling"
    p5t.alignment = PP_ALIGN.LEFT
    run5t = p5t.runs[0]
    run5t.font.size = Pt(32)
    run5t.font.bold = True
    run5t.font.color.rgb = RGBColor(0x1F, 0x48, 0x7E)

    # Body textbox - LEFT aligned (will remain LEFT after task)
    body5_width = Inches(7)
    body5_left = (slide_width - body5_width) // 2
    body5 = slide5.shapes.add_textbox(body5_left, Inches(1.4), body5_width, Inches(4.5))
    tf5 = body5.text_frame
    tf5.word_wrap = True

    lines5 = [
        "Linear and logistic regression fundamentals.",
        "Model evaluation: accuracy, precision, recall, F1-score.",
        "Cross-validation strategies to prevent overfitting.",
        "Feature selection and dimensionality reduction techniques.",
        "Deploying models using scikit-learn pipelines.",
    ]
    for i, line in enumerate(lines5):
        if i == 0:
            p = tf5.paragraphs[0]
        else:
            p = tf5.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 6: Module 4 ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])

    title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.9))
    tf6t = title6.text_frame
    p6t = tf6t.paragraphs[0]
    p6t.text = "Module 4: Machine Learning Applications"
    p6t.alignment = PP_ALIGN.LEFT
    run6t = p6t.runs[0]
    run6t.font.size = Pt(32)
    run6t.font.bold = True
    run6t.font.color.rgb = RGBColor(0x1F, 0x48, 0x7E)

    # Body textbox - LEFT aligned (will be changed to RIGHT)
    body6_width = Inches(7)
    body6_left = (slide_width - body6_width) // 2
    body6 = slide6.shapes.add_textbox(body6_left, Inches(1.4), body6_width, Inches(4.5))
    tf6 = body6.text_frame
    tf6.word_wrap = True

    lines6 = [
        "Supervised learning: classification and regression trees.",
        "Unsupervised learning: k-means and hierarchical clustering.",
        "Ensemble methods: random forests and gradient boosting.",
        "Neural network basics and deep learning introduction.",
        "Practical project: customer segmentation analysis.",
    ]
    for i, line in enumerate(lines6):
        if i == 0:
            p = tf6.paragraphs[0]
        else:
            p = tf6.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 7: Closing / Summary ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])

    title7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.9))
    tf7t = title7.text_frame
    p7t = tf7t.paragraphs[0]
    p7t.text = "Workshop Summary & Next Steps"
    p7t.alignment = PP_ALIGN.LEFT
    run7t = p7t.runs[0]
    run7t.font.size = Pt(32)
    run7t.font.bold = True
    run7t.font.color.rgb = RGBColor(0x1F, 0x48, 0x7E)

    body7_width = Inches(7)
    body7_left = (slide_width - body7_width) // 2
    body7 = slide7.shapes.add_textbox(body7_left, Inches(1.4), body7_width, Inches(4.5))
    tf7 = body7.text_frame
    tf7.word_wrap = True

    lines7 = [
        "Complete all four modules for certification.",
        "Submit capstone project by April 15, 2025.",
        "Join our analytics community forum for ongoing support.",
        "Optional advanced track available from May 2025.",
        "Contact workshop@dataanalytics.org for questions.",
    ]
    for i, line in enumerate(lines7):
        if i == 0:
            p = tf7.paragraphs[0]
        else:
            p = tf7.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
