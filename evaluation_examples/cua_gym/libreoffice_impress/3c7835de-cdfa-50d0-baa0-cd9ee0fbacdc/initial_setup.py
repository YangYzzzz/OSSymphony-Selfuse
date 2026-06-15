"""
Initial Setup: Create presentation with blank slide 3 for timeline task
Task ID: impress_gf5_008
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_008'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Executive Summary Presentation"
    slide1.placeholders[1].text = "Q4 2024 Strategic Review\nAcme Corporation"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Financial Performance Overview"
    items = [
        "Product Launch Timeline",
        "Market Expansion Strategy",
        "Key Milestones & Achievements",
        "Next Steps & Action Items",
    ]
    for item in items:
        p = body2.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 3: Project Timeline (BLANK - no drawing objects) ---
    # Use Title Only layout (index 5 = Blank, but we want the title)
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    # Add a title text box manually at the top
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Project Timeline"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x3B, 0x4E)
    # No other shapes on slide 3 - the agent must draw the timeline

    # --- Slide 4: Key Milestones ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Key Milestones"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "2021: Foundation Phase - Core platform development completed"
    milestones = [
        "2022: Growth Phase - Expanded to 3 new markets",
        "2023: Scale Phase - Revenue crossed $10M ARR",
        "2024: Maturity Phase - IPO preparation underway",
    ]
    for m in milestones:
        p = body4.add_paragraph()
        p.text = m
        p.level = 0

    # --- Slide 5: Next Steps ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Next Steps"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Finalize Series C funding round by Q1 2025"
    next_items = [
        "Launch product v3.0 in European markets",
        "Hire 50 additional engineers for platform team",
        "Complete SOC2 Type II certification",
    ]
    for item in next_items:
        p = body5.add_paragraph()
        p.text = item
        p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
