"""
Reward Script: Timeline on slide 3 with arrow, tick marks, year labels, grouped
Task ID: impress_gf5_008
Domain: libreoffice_impress
Scoring:
  Component 1: Group shape exists on slide 3 (0.2)
  Component 2: Horizontal arrow spanning ~80% width (0.2)
  Component 3: 4 vertical tick marks (0.2)
  Component 4: Year labels 2021-2024 present (0.2)
  Component 5: Labels positioned below the arrow/ticks (0.2)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_008'


def persist_app_state(domain: str):
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_all_shapes_in_group(group_shape):
    """Recursively collect all shapes from a group."""
    shapes = []
    if hasattr(group_shape, 'shapes'):
        for sub in group_shape.shapes:
            shapes.append(sub)
            shapes.extend(get_all_shapes_in_group(sub))
    return shapes


def verify_task(file_path):
    """
    Verify timeline creation on slide 3 with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: at least 3 slides exist
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[2]  # slide 3 (0-indexed)
    slide_width = prs.slide_width

    # Find group shape(s) on slide 3
    group_shapes = []
    for shape in slide.shapes:
        if hasattr(shape, 'shapes') and len(list(shape.shapes)) > 0:
            group_shapes.append(shape)

    # Component 1: Group shape exists on slide 3 (0.2 points)
    # This check distinguishes golden (has group) from initial (no group)
    try:
        if len(group_shapes) > 0:
            print(f"PASS: Component 1 — Group shape found on slide 3 ({len(group_shapes)} group(s)) (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 1 — No group shape found on slide 3")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    if len(group_shapes) == 0:
        # No group means no timeline — skip remaining checks
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Use the largest group (most sub-shapes) as the timeline group
    timeline_group = max(group_shapes, key=lambda g: len(list(g.shapes)))
    sub_shapes = get_all_shapes_in_group(timeline_group)

    # Component 2: Horizontal arrow/line spanning ~80% of slide width (0.2 points)
    # Look for a wide horizontal shape (width >= 60% of slide width, effectively zero height)
    try:
        arrow_candidates = [
            sub for sub in sub_shapes
            if sub.width >= slide_width * 0.6 and sub.height <= slide_width * 0.1
        ]
        if len(arrow_candidates) > 0:
            arrow_ratio = arrow_candidates[0].width / slide_width * 100
            if arrow_ratio >= 60.0:
                print(f"PASS: Component 2 — Horizontal arrow found, width spans {arrow_ratio:.1f}% of slide (0.2 pts)")
                total_score += 0.2
        else:
            widths = [sub.width for sub in sub_shapes]
            print(f"FAIL: Component 2 — No horizontal arrow found. Sub-shape widths: {widths}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: 4 vertical tick marks (0.2 points)
    # Tick marks are shapes with zero (or near-zero) width and non-zero height
    try:
        tick_shapes = []
        for sub in sub_shapes:
            # Tick: narrow width, some height, not a text box with text
            is_text_with_content = hasattr(sub, 'text') and sub.text.strip() != ''
            if not is_text_with_content and sub.width <= slide_width * 0.02 and sub.height > 0:
                tick_shapes.append(sub)

        if len(tick_shapes) >= 4:
            print(f"PASS: Component 3 — {len(tick_shapes)} vertical tick marks found (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 3 — Expected 4 tick marks, found {len(tick_shapes)}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Year labels 2021-2024 present in group (0.2 points)
    try:
        required_labels = {'2021', '2022', '2023', '2024'}
        found_labels = set()
        label_shapes = []
        for sub in sub_shapes:
            if hasattr(sub, 'text_frame'):
                text = sub.text_frame.text.strip()
                if text in required_labels:
                    found_labels.add(text)
                    label_shapes.append(sub)

        if found_labels == required_labels:
            print(f"PASS: Component 4 — All 4 year labels found: {sorted(found_labels)} (0.2 pts)")
            total_score += 0.2
        else:
            missing = required_labels - found_labels
            print(f"FAIL: Component 4 — Missing labels: {sorted(missing)}, found: {sorted(found_labels)}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Labels positioned below the arrow/ticks (0.2 points)
    # The label tops should be below the tick mark tops
    try:
        if len(tick_shapes) >= 4 and len(label_shapes) >= 4:
            # Average top of ticks vs average top of labels
            avg_tick_top = sum(t.top for t in tick_shapes[:4]) / 4
            avg_label_top = sum(l.top for l in label_shapes[:4]) / 4

            if avg_label_top > avg_tick_top:
                print(f"PASS: Component 5 — Labels (avg top={avg_label_top}) below ticks (avg top={avg_tick_top}) (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 5 — Labels (avg top={avg_label_top}) not below ticks (avg top={avg_tick_top})")
        else:
            print(f"FAIL: Component 5 — Insufficient tick/label shapes to compare positions")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
