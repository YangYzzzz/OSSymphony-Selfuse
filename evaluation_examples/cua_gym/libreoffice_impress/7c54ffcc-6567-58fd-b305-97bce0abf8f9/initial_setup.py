"""
Initial Setup: Reordering overlapping shapes in LibreOffice Impress
Task ID: impress_objects_040
Domain: libreoffice_impress

Creates a PPTX file with:
- Slide 1: Title slide
- Slide 2: 3 overlapping shapes where red triangle is at the back
  (stacking order back->front: red triangle, green rectangle, blue circle)
- Slide 3: Content slide
The file is placed at ~/Desktop/layers_test.pptx
"""

import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_objects_040'
OUTPUT = f'{WORKDIR}/{TASK_ID}_initial.pptx'
DESKTOP_OUTPUT = f'{WORKDIR}/Desktop/layers_test.pptx'


def add_shape_with_color(slide, shape_type, left, top, width, height, fill_color, name):
    """Add an auto shape with solid fill color."""
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.name = name
    # Set fill color
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_color
    # Remove line or set to matching color
    shape.line.color.rgb = fill_color
    return shape


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Shape Layering Demo"
    slide1.placeholders[1].text = "A demonstration of overlapping shapes and Z-order"

    # ---- Slide 2: Overlapping Shapes (the task slide) ----
    # Stacking order back->front: red triangle, green rectangle, blue circle
    # In python-pptx, shapes added later appear on top
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Add title text box
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.paragraphs[0].text = "Overlapping Shapes"
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Center of the slide for overlapping
    cx = Inches(5)
    cy = Inches(4)
    shape_size = Inches(2.5)

    # Overlap offset
    offset = Inches(0.6)

    # Shape 1 (BACK): Red Triangle
    # MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE = 5
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    red_triangle = slide2.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        cx - shape_size / 2,
        cy - shape_size / 2,
        shape_size,
        shape_size
    )
    red_triangle.name = "Red Triangle"
    red_triangle.fill.solid()
    red_triangle.fill.fore_color.rgb = RGBColor(0xFF, 0x00, 0x00)
    red_triangle.line.color.rgb = RGBColor(0xCC, 0x00, 0x00)
    red_triangle.line.width = Pt(1.5)

    # Shape 2 (MIDDLE): Green Rectangle
    green_rect = slide2.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        cx - shape_size / 2 + offset,
        cy - shape_size / 2 + offset,
        shape_size,
        shape_size
    )
    green_rect.name = "Green Rectangle"
    green_rect.fill.solid()
    green_rect.fill.fore_color.rgb = RGBColor(0x00, 0xB0, 0x50)
    green_rect.line.color.rgb = RGBColor(0x00, 0x80, 0x30)
    green_rect.line.width = Pt(1.5)

    # Shape 3 (FRONT): Blue Circle (Oval)
    blue_circle = slide2.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        cx - shape_size / 2 - offset,
        cy - shape_size / 2 + offset,
        shape_size,
        shape_size
    )
    blue_circle.name = "Blue Circle"
    blue_circle.fill.solid()
    blue_circle.fill.fore_color.rgb = RGBColor(0x00, 0x70, 0xC0)
    blue_circle.line.color.rgb = RGBColor(0x00, 0x50, 0xA0)
    blue_circle.line.width = Pt(1.5)

    # Add labels inside each shape (use add_run for proper font control)
    def set_shape_label(shape, text, color_rgb):
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        for run in p.runs:
            run.font.color.rgb = color_rgb
            run.font.size = Pt(14)
            run.font.bold = True

    set_shape_label(red_triangle, "Red Triangle", RGBColor(0xFF, 0xFF, 0xFF))
    set_shape_label(green_rect, "Green Rectangle", RGBColor(0xFF, 0xFF, 0xFF))
    set_shape_label(blue_circle, "Blue Circle", RGBColor(0xFF, 0xFF, 0xFF))

    # ---- Slide 3: Content/Notes Slide ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide3.shapes.title.text = "About Shape Layering"
    content_tf = slide3.placeholders[1].text_frame
    content_tf.text = "Z-order controls which shapes appear in front"
    p2 = content_tf.add_paragraph()
    p2.text = "Right-click a shape to access 'Send to Back' and 'Bring to Front'"
    p2.level = 1
    p3 = content_tf.add_paragraph()
    p3.text = "The stacking order is stored in the XML element sequence"
    p3.level = 1
    p4 = content_tf.add_paragraph()
    p4.text = "Use 'Bring Forward' and 'Send Backward' for fine control"
    p4.level = 1

    # Save
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Also copy to Desktop as layers_test.pptx (as required by task context)
    os.makedirs(f'{WORKDIR}/Desktop', exist_ok=True)
    shutil.copy(OUTPUT, DESKTOP_OUTPUT)
    print(f'Also copied to: {DESKTOP_OUTPUT}')

    # Verify stacking order
    verify = Presentation(OUTPUT)
    slide_shapes = verify.slides[1].shapes  # slide 2 (0-indexed = 1)
    print(f'\nInitial stacking order on Slide 2 (back to front):')
    for i, s in enumerate(slide_shapes):
        print(f'  [{i}] {s.name}')


create_initial()
