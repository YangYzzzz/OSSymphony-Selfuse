"""
Reward Script: Reorder overlapping shapes on slide 2 of layers_test.pptx
Task ID: impress_objects_040
Domain: libreoffice_impress
Scoring:
  - Component 1: Blue Circle is behind (lower z-order than) Green Rectangle — 0.35 pts
  - Component 2: Green Rectangle is behind (lower z-order than) Red Triangle — 0.30 pts
  - Component 3: Blue Circle is behind (lower z-order than) Red Triangle — 0.35 pts
  Total: 1.0

Z-order in python-pptx: position in slide.shapes list determines stacking.
Lower index = further back (behind other shapes), higher index = closer to front.
Expected back-to-front ordering: Blue Circle (0070C0), Green Rectangle (00B050), Red Triangle (FF0000).

Using pairwise ordering checks ensures initial file scores 0.0:
  Initial order (back-to-front): Red Triangle, Green Rectangle, Blue Circle
  - Blue Circle behind Green Rectangle? NO (Blue Circle index 4 > Green Rectangle index 3)
  - Green Rectangle behind Red Triangle? NO (Green Rectangle index 3 > Red Triangle index 2)
  - Blue Circle behind Red Triangle? NO (Blue Circle index 4 > Red Triangle index 2)
  → All fail → initial score = 0.0

  Golden order (back-to-front): Blue Circle, Green Rectangle, Red Triangle
  - Blue Circle behind Green Rectangle? YES (Blue Circle index 2 < Green Rectangle index 3)
  - Green Rectangle behind Red Triangle? YES (Green Rectangle index 3 < Red Triangle index 4)
  - Blue Circle behind Red Triangle? YES (Blue Circle index 2 < Red Triangle index 4)
  → All pass → golden score = 1.0
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'impress_objects_040'

# Color constants (uppercase hex RGB as returned by python-pptx str(rgb).upper())
BLUE_CIRCLE_COLOR = '0070C0'
GREEN_RECT_COLOR = '00B050'
RED_TRIANGLE_COLOR = 'FF0000'


def get_shape_color(shape):
    """Return uppercase hex RGB string if shape has solid fill, else None."""
    try:
        if shape.fill.type == 1:  # MSO_FILL.SOLID
            return str(shape.fill.fore_color.rgb).upper()
    except Exception:
        pass
    return None


def verify_task(file_path):
    """
    Verify z-order of 3 overlapping shapes on slide 2.
    Expected back-to-front: Blue Circle, Green Rectangle, Red Triangle.
    Uses pairwise ordering checks (not absolute position checks) so that
    partial credit is only awarded for genuinely-changed orderings.
    Returns float between 0.0 and 1.0.
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"CRITICAL: Expected at least 2 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[1]  # slide 2 (0-indexed)

    # Collect shapes and their z-order indices (position in shapes list)
    # Lower index = further back in the stacking order
    shape_z_order = {}  # color_hex -> z_order_index
    shape_names = {}    # color_hex -> shape name (for logging)
    for i, shape in enumerate(slide.shapes):
        color = get_shape_color(shape)
        if color in (BLUE_CIRCLE_COLOR, GREEN_RECT_COLOR, RED_TRIANGLE_COLOR):
            shape_z_order[color] = i
            shape_names[color] = shape.name

    if len(shape_z_order) != 3:
        print(f"CRITICAL: Expected 3 colored shapes on slide 2, found {len(shape_z_order)}: {shape_z_order}")
        print("REWARD: 0.0")
        return 0.0

    blue_z = shape_z_order[BLUE_CIRCLE_COLOR]
    green_z = shape_z_order[GREEN_RECT_COLOR]
    red_z = shape_z_order[RED_TRIANGLE_COLOR]

    print(f"Z-order indices (lower = further back):")
    print(f"  Blue Circle ({shape_names[BLUE_CIRCLE_COLOR]}): z={blue_z}")
    print(f"  Green Rectangle ({shape_names[GREEN_RECT_COLOR]}): z={green_z}")
    print(f"  Red Triangle ({shape_names[RED_TRIANGLE_COLOR]}): z={red_z}")

    # Component 1: Blue Circle is behind (lower z-order than) Green Rectangle (0.35 points)
    # This FAILS in initial (blue_z=4, green_z=3 → blue is IN FRONT) → correct
    # This PASSES in golden (blue_z=2, green_z=3 → blue is BEHIND) → correct
    try:
        if blue_z < green_z:
            print(f"PASS: Component 1 — Blue Circle (z={blue_z}) is behind Green Rectangle (z={green_z}) (0.35 pts)")
            total_score += 0.35
        else:
            print(f"FAIL: Component 1 — Blue Circle (z={blue_z}) is NOT behind Green Rectangle (z={green_z}); expected blue_z < green_z")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Green Rectangle is behind (lower z-order than) Red Triangle (0.30 points)
    # This FAILS in initial (green_z=3, red_z=2 → green is IN FRONT) → correct
    # This PASSES in golden (green_z=3, red_z=4 → green is BEHIND) → correct
    try:
        if green_z < red_z:
            print(f"PASS: Component 2 — Green Rectangle (z={green_z}) is behind Red Triangle (z={red_z}) (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 2 — Green Rectangle (z={green_z}) is NOT behind Red Triangle (z={red_z}); expected green_z < red_z")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Blue Circle is behind (lower z-order than) Red Triangle (0.35 points)
    # This FAILS in initial (blue_z=4, red_z=2 → blue is IN FRONT) → correct
    # This PASSES in golden (blue_z=2, red_z=4 → blue is BEHIND) → correct
    try:
        if blue_z < red_z:
            print(f"PASS: Component 3 — Blue Circle (z={blue_z}) is behind Red Triangle (z={red_z}) (0.35 pts)")
            total_score += 0.35
        else:
            print(f"FAIL: Component 3 — Blue Circle (z={blue_z}) is NOT behind Red Triangle (z={red_z}); expected blue_z < red_z")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against the Desktop file (the file the agent will have modified)
file_path = '/home/user/Desktop/layers_test.pptx'
if not os.path.exists(file_path):
    # Fallback to golden file for testing
    file_path = f'{WORKDIR}/{TASK_ID}_initial.pptx'
    if not os.path.exists(file_path):
        print(f"File not found: {file_path}")
        print("REWARD: 0.0")
    else:
        verify_task(file_path)
else:
    verify_task(file_path)
