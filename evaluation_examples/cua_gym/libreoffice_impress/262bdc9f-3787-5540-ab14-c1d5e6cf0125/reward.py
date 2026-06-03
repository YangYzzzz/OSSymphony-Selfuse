"""
Reward Script: Scientific poster slide at 90x120cm portrait with structured layout
Task ID: impress_gf2_027
Domain: libreoffice_impress
Scoring:
  C1: Slide dimensions 90x120cm portrait (0.20)
  C2: Title rectangle — full width, 15cm tall, #B91C1C fill (0.20)
  C3: Authors rectangle — full width, 8cm tall, #F3F4F6 fill (0.20)
  C4: Three body columns (~28.7cm wide each) for Intro/Methods/Results (0.25)
  C5: Title text is white, columns have heading text (0.15)
"""

import os
import sys

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_027'

# Tolerances
CM_TOL = 2.0        # cm tolerance for dimensions
CM_TO_EMU = 360000   # 1 cm = 360000 EMU
REL_TOL = 0.10       # 10% relative tolerance for shape sizes


def approx_cm(emu_val, expected_cm, tol_cm=CM_TOL):
    """Check if EMU value is approximately expected_cm within tolerance."""
    actual_cm = emu_val / CM_TO_EMU
    return abs(actual_cm - expected_cm) <= tol_cm


def get_fill_rgb(shape):
    """Get fill color RGB string from a shape, or None."""
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def get_all_text(shape):
    """Get concatenated text from all runs in a shape."""
    if not shape.has_text_frame:
        return ""
    texts = []
    for para in shape.text_frame.paragraphs:
        texts.append(para.text)
    return " ".join(texts).strip()


def get_text_color_rgb(shape):
    """Get the RGB color of the first non-empty run's font, or None."""
    if not shape.has_text_frame:
        return None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if (run.text or "").strip():
                try:
                    if run.font.color.type is not None:
                        return str(run.font.color.rgb)
                except Exception:
                    pass
    return None


def verify_task(file_path):
    """Verify task completion with progressive scoring. Returns float 0.0-1.0."""
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    if len(prs.slides) == 0:
        print("FAIL: No slides found")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]
    shapes = list(slide.shapes)

    # =========================================================
    # Component 1: Slide dimensions 90x120cm portrait (0.20 pts)
    # =========================================================
    try:
        w_cm = prs.slide_width / CM_TO_EMU
        h_cm = prs.slide_height / CM_TO_EMU
        is_portrait = prs.slide_width < prs.slide_height
        dim_ok = approx_cm(prs.slide_width, 90.0, 3.0) and approx_cm(prs.slide_height, 120.0, 3.0)

        if is_portrait and dim_ok:
            print(f"PASS: Component 1 — Slide is {w_cm:.1f}x{h_cm:.1f} cm portrait (0.20 pts)")
            total_score += 0.20
        elif is_portrait:
            # Partial: portrait but wrong dimensions
            print(f"FAIL: Component 1 — Portrait but dimensions {w_cm:.1f}x{h_cm:.1f} cm, expected ~90x120")
        elif dim_ok:
            print(f"FAIL: Component 1 — Dimensions correct but not portrait")
        else:
            print(f"FAIL: Component 1 — Slide is {w_cm:.1f}x{h_cm:.1f} cm landscape, expected 90x120 portrait")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # =========================================================
    # Component 2: Title rectangle — full width, ~15cm tall, #B91C1C fill (0.20 pts)
    # =========================================================
    try:
        title_rect = None
        for shape in shapes:
            fill_rgb = get_fill_rgb(shape)
            if fill_rgb and fill_rgb.upper() == 'B91C1C':
                title_rect = shape
                break

        if title_rect is not None:
            tr_w = title_rect.width / CM_TO_EMU
            tr_h = title_rect.height / CM_TO_EMU
            tr_top = title_rect.top / CM_TO_EMU
            # Full width: at least 80% of slide width
            width_ok = title_rect.width >= prs.slide_width * 0.80
            height_ok = approx_cm(title_rect.height, 15.0, 4.0)
            top_ok = tr_top < 5.0  # near top

            if width_ok and height_ok and top_ok:
                print(f"PASS: Component 2 — Title rect {tr_w:.1f}x{tr_h:.1f} cm at top with #B91C1C (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 2 — Title rect found (#B91C1C) but dims off: {tr_w:.1f}x{tr_h:.1f} cm, top={tr_top:.1f}")
        else:
            print("FAIL: Component 2 — No rectangle with #B91C1C fill found")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # =========================================================
    # Component 3: Authors rectangle — full width, ~8cm tall, #F3F4F6 fill (0.20 pts)
    # =========================================================
    try:
        authors_rect = None
        for shape in shapes:
            fill_rgb = get_fill_rgb(shape)
            if fill_rgb and fill_rgb.upper() == 'F3F4F6':
                authors_rect = shape
                break

        if authors_rect is not None:
            ar_w = authors_rect.width / CM_TO_EMU
            ar_h = authors_rect.height / CM_TO_EMU
            ar_top = authors_rect.top / CM_TO_EMU
            width_ok = authors_rect.width >= prs.slide_width * 0.80
            height_ok = approx_cm(authors_rect.height, 8.0, 3.0)
            # Should be below title area (roughly at 15cm or so)
            position_ok = 10.0 < ar_top < 30.0

            if width_ok and height_ok and position_ok:
                print(f"PASS: Component 3 — Authors rect {ar_w:.1f}x{ar_h:.1f} cm at top={ar_top:.1f} cm with #F3F4F6 (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 3 — Authors rect found (#F3F4F6) but dims/pos off: {ar_w:.1f}x{ar_h:.1f} cm, top={ar_top:.1f}")
        else:
            print("FAIL: Component 3 — No rectangle with #F3F4F6 fill found")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # =========================================================
    # Component 4: Three body columns for Intro/Methods/Results (0.25 pts)
    # =========================================================
    try:
        # Look for column-like shapes: tall, narrow-ish, positioned below authors section
        column_shapes = []
        for shape in shapes:
            s_top = shape.top / CM_TO_EMU
            s_h = shape.height / CM_TO_EMU
            s_w = shape.width / CM_TO_EMU
            # Column criteria: below the authors area (top > 15cm), tall (> 30cm), width ~25-35cm
            if s_top > 15.0 and s_h > 30.0 and 15.0 < s_w < 45.0:
                column_shapes.append(shape)

        # Check that we have 3 columns
        if len(column_shapes) >= 3:
            # Sort by left position
            column_shapes.sort(key=lambda s: s.left)
            # Check they contain the right headings
            col_texts = [get_all_text(s).lower() for s in column_shapes[:3]]
            has_intro = any('introduction' in t for t in col_texts)
            has_methods = any('method' in t for t in col_texts)
            has_results = any('result' in t for t in col_texts)

            if has_intro and has_methods and has_results:
                widths = [s.width / CM_TO_EMU for s in column_shapes[:3]]
                print(f"PASS: Component 4 — 3 columns found: Intro/Methods/Results, widths={[f'{w:.1f}' for w in widths]} cm (0.25 pts)")
                total_score += 0.25
            elif len(column_shapes) >= 3:
                # Partial: 3 columns but missing labels
                print(f"PARTIAL: Component 4 — 3 column shapes found but labels missing/wrong: {col_texts} (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 4 — Found {len(column_shapes)} column-like shapes, need 3")
        else:
            # Maybe columns are smaller or differently shaped; check for any 3 shapes below authors
            below_shapes = [s for s in shapes if s.top / CM_TO_EMU > 15.0]
            if len(below_shapes) >= 3:
                print(f"PARTIAL: Component 4 — {len(below_shapes)} shapes below authors but don't meet column criteria (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 4 — Only {len(column_shapes)} column-like shapes found, need 3")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # =========================================================
    # Component 5: Title text is white (0.15 pts)
    # =========================================================
    try:
        if title_rect is not None:
            title_text = get_all_text(title_rect)
            title_color = get_text_color_rgb(title_rect)

            if title_text and title_color and title_color.upper() == 'FFFFFF':
                print(f"PASS: Component 5 — Title has white text: \"{title_text}\" (0.15 pts)")
                total_score += 0.15
            elif title_text and title_color:
                print(f"FAIL: Component 5 — Title text color is #{title_color}, expected #FFFFFF")
            elif title_text:
                print(f"FAIL: Component 5 — Title has text but color could not be read")
            else:
                print(f"FAIL: Component 5 — Title rectangle has no text")
        else:
            print("FAIL: Component 5 — No title rectangle found (depends on C2)")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
