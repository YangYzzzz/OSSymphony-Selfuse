"""
Initial Setup: Create a blank presentation with default widescreen size
Task ID: impress_gf2_027
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_027'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Default widescreen slide (25.4cm x 19.05cm landscape)
    # Add one blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only layout
    # Clear any placeholder text
    for shape in slide.placeholders:
        sp = shape._element
        sp.getparent().remove(sp)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
