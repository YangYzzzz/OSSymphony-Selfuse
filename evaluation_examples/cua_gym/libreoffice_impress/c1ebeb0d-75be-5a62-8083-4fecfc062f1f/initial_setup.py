"""
Initial Setup: Create a 7-slide sales pitch presentation with slide 4 having only a title.
Task ID: impress_sales_071
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_071'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_text(slide, text, font_size=Pt(32), bold=True, color=RGBColor(0x1A, 0x1A, 0x2E)):
    """Set title text with formatting on a slide."""
    title = slide.shapes.title
    title.text = text
    for run in title.text_frame.paragraphs[0].runs:
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color


def add_body_text(slide, placeholder_idx, lines, font_size=Pt(18)):
    """Add bullet-point text to a content placeholder."""
    ph = slide.placeholders[placeholder_idx]
    tf = ph.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.level = 0
        for run in p.runs:
            run.font.size = font_size
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "FeatureList Pitch"
    slide1.placeholders[1].text = "SyncWave Platform — Enterprise Product Overview"
    for run in slide1.shapes.title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    for run in slide1.placeholders[1].text_frame.paragraphs[0].runs:
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0x4A, 0x5A, 0x6A)

    # --- Slide 2: Market Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide2, "Market Overview")
    add_body_text(slide2, 1, [
        "Enterprise SaaS market projected to reach $232B by 2027",
        "68% of companies plan to increase cloud spending this year",
        "Data synchronization remains top IT challenge for 45% of CIOs",
        "Average enterprise uses 110+ SaaS applications",
        "Integration and interoperability drive purchasing decisions",
    ])

    # --- Slide 3: Target Audience ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide3, "Target Audience")
    add_body_text(slide3, 1, [
        "Mid-market companies (500-5000 employees)",
        "IT Directors and VP of Engineering",
        "Operations teams managing cross-platform workflows",
        "Compliance officers in regulated industries",
        "CTOs evaluating digital transformation initiatives",
    ])

    # --- Slide 4: Core Features (TITLE ONLY — no table, no content) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Core Features"
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # --- Slide 5: Pricing ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide5, "Pricing Plans")
    add_body_text(slide5, 1, [
        "Starter: $29/user/month — Up to 25 users, 10GB storage",
        "Professional: $59/user/month — Unlimited users, 100GB storage",
        "Enterprise: Custom pricing — Dedicated infrastructure, SLA",
        "All plans include 14-day free trial",
        "Annual billing saves 20% across all tiers",
    ])

    # --- Slide 6: Product Roadmap ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide6, "Product Roadmap")
    add_body_text(slide6, 1, [
        "Q2 2026: Enhanced analytics dashboard with custom KPIs",
        "Q3 2026: Native mobile apps for iOS and Android",
        "Q4 2026: Advanced workflow automation engine",
        "Q1 2027: Multi-region data residency options",
        "Q2 2027: AI-powered anomaly detection and alerts",
    ])

    # --- Slide 7: Contact Us ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide7, "Contact Us")
    add_body_text(slide7, 1, [
        "Website: www.syncwave.io",
        "Sales: sales@syncwave.io | +1 (415) 555-0192",
        "Support: support@syncwave.io",
        "Schedule a demo: demo.syncwave.io",
        "Follow us: @SyncWaveHQ on LinkedIn and Twitter",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
