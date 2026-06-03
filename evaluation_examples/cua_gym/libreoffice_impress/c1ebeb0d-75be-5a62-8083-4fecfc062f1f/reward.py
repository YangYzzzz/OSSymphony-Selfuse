"""
Reward Script: Insert a 5-row by 2-column feature table on slide 4
Task ID: impress_sales_071
Domain: libreoffice_impress
Scoring:
  C1 (0.20) - Table exists on slide 4 with 6 rows x 2 cols
  C2 (0.15) - Header row content correct ("Feature", "Description")
  C3 (0.30) - Data rows content matches (5 feature/description pairs)
  C4 (0.20) - Header text bold + white (FFFFFF)
  C5 (0.15) - Header cell background #2B6CB0
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_071'

# Expected table data (row 0 = header, rows 1-5 = data)
EXPECTED_HEADERS = ['Feature', 'Description']
EXPECTED_DATA = [
    ('Auto-Sync', 'Real-time data synchronization across all devices'),
    ('Smart Analytics', 'AI-powered insights and recommendations'),
    ('Team Hub', 'Centralized collaboration workspace'),
    ('API Gateway', 'RESTful API with 99.9% uptime'),
    ('Compliance Suite', 'SOC2 and GDPR built-in'),
]


def persist_app_state():
    """Try to save any unsaved LibreOffice state."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide4 = prs.slides[3]  # 0-indexed

    # Find table shape on slide 4
    table_shape = None
    for shape in slide4.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_shape = shape
            break

    if table_shape is None:
        print("FAIL: No table found on slide 4")
        print("REWARD: 0.0")
        return 0.0

    table = table_shape.table

    # Component 1: Table dimensions (0.20 points)
    try:
        num_rows = len(table.rows)
        num_cols = len(table.columns)
        if num_rows == 6 and num_cols == 2:
            print(f"PASS: Component 1 — Table is 6x2 (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 — Table is {num_rows}x{num_cols}, expected 6x2")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Header row content (0.15 points)
    try:
        h0 = table.cell(0, 0).text.strip()
        h1 = table.cell(0, 1).text.strip()
        if h0 == EXPECTED_HEADERS[0] and h1 == EXPECTED_HEADERS[1]:
            print(f"PASS: Component 2 — Headers correct: '{h0}', '{h1}' (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 — Headers: '{h0}', '{h1}'; expected '{EXPECTED_HEADERS[0]}', '{EXPECTED_HEADERS[1]}'")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Data rows content (0.30 points — 0.06 per row)
    try:
        data_score = 0.0
        for i, (exp_feat, exp_desc) in enumerate(EXPECTED_DATA):
            row_idx = i + 1
            try:
                actual_feat = table.cell(row_idx, 0).text.strip()
                actual_desc = table.cell(row_idx, 1).text.strip()
                if actual_feat == exp_feat and actual_desc == exp_desc:
                    data_score += 0.06
                    print(f"  Row {row_idx}: PASS — '{actual_feat}' / '{actual_desc}'")
                else:
                    print(f"  Row {row_idx}: FAIL — got '{actual_feat}' / '{actual_desc}', expected '{exp_feat}' / '{exp_desc}'")
            except Exception as e:
                print(f"  Row {row_idx}: ERROR — {e}")
        if data_score > 0:
            total_score += data_score
        if data_score >= 0.29:
            print(f"PASS: Component 3 — All data rows correct ({data_score:.2f} pts)")
        else:
            print(f"PARTIAL: Component 3 — Data rows ({data_score:.2f}/0.30 pts)")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Header text bold + white (0.20 points)
    try:
        header_style_score = 0.0
        for col_idx in range(2):
            cell = table.cell(0, col_idx)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    is_bold = run.font.bold is True
                    try:
                        is_white = str(run.font.color.rgb).upper() == 'FFFFFF'
                    except Exception:
                        is_white = False
                    if is_bold and is_white:
                        header_style_score += 0.10
                        print(f"  Header [{0},{col_idx}]: PASS — bold={is_bold}, color=white")
                    else:
                        print(f"  Header [{0},{col_idx}]: FAIL — bold={run.font.bold}, white={is_white}")
                    break  # check first run only
                break  # check first paragraph only
        if header_style_score > 0:
            total_score += header_style_score
        if header_style_score >= 0.19:
            print(f"PASS: Component 4 — Header text bold+white ({header_style_score:.2f} pts)")
        else:
            print(f"PARTIAL: Component 4 — Header text styling ({header_style_score:.2f}/0.20 pts)")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Header cell background #2B6CB0 (0.15 points)
    try:
        bg_score = 0.0
        for col_idx in range(2):
            cell = table.cell(0, col_idx)
            try:
                fill = cell.fill
                if fill.type is not None and fill.type == 1:  # SOLID
                    color_hex = str(fill.fore_color.rgb).upper()
                    if color_hex == '2B6CB0':
                        bg_score += 0.075
                        print(f"  Header bg [{0},{col_idx}]: PASS — {color_hex}")
                    else:
                        print(f"  Header bg [{0},{col_idx}]: FAIL — got {color_hex}, expected 2B6CB0")
                else:
                    print(f"  Header bg [{0},{col_idx}]: FAIL — fill type is {fill.type}, not SOLID")
            except Exception as e:
                print(f"  Header bg [{0},{col_idx}]: ERROR — {e}")
        if bg_score > 0:
            total_score += bg_score
        if bg_score >= 0.14:
            print(f"PASS: Component 5 — Header background #2B6CB0 ({bg_score:.3f} pts)")
        else:
            print(f"PARTIAL: Component 5 — Header background ({bg_score:.3f}/0.15 pts)")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
