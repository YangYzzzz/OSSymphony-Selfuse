"""
Initial Setup: Dark-themed presentation master design task
Task ID: impress_gf4_016
Domain: libreoffice_impress

Creates a 14-slide Tech Conference presentation with plain white background
and black default fonts. No dark theme, no gradient, no colored line.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_016'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_content(slide, title_text, body_lines):
    """Helper to populate a slide with title and body content."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
            run.font.size = Pt(32)
            run.font.name = "Calibri"

    # Find body placeholder (index 1 typically)
    body_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            body_ph = ph
            break
    if body_ph and body_lines:
        tf = body_ph.text_frame
        tf.clear()
        for i, line in enumerate(body_lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            for run in p.runs:
                run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
                run.font.size = Pt(18)
                run.font.name = "Calibri"


def create_initial():
    prs = Presentation()

    # Standard 16:9 dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide layouts from default template
    title_layout = prs.slide_layouts[0]       # Title Slide
    content_layout = prs.slide_layouts[1]     # Title + Content
    blank_layout = prs.slide_layouts[6]       # Title Only

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(title_layout)
    slide1.shapes.title.text = "Tech Conference 2026"
    if slide1.placeholders[1]:
        slide1.placeholders[1].text = "Innovation in the Age of AI"
    for ph in slide1.placeholders:
        for p in ph.text_frame.paragraphs:
            for run in p.runs:
                run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
                run.font.name = "Calibri"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(content_layout)
    add_title_content(slide2, "Conference Agenda", [
        "09:00 - Registration & Networking Breakfast",
        "10:00 - Keynote: The Future of Distributed Computing",
        "11:30 - Panel: AI Ethics in Enterprise Software",
        "13:00 - Lunch Break & Demo Stations",
        "14:30 - Workshop Track A: Cloud-Native Architecture",
        "14:30 - Workshop Track B: Machine Learning Pipelines",
        "16:00 - Lightning Talks & Community Showcase",
        "17:30 - Closing Remarks & Evening Reception",
    ])

    # --- Slide 3: Keynote Speaker ---
    slide3 = prs.slides.add_slide(content_layout)
    add_title_content(slide3, "Keynote: Dr. Amara Okafor", [
        "Chief Technology Officer, NexusAI Systems",
        "20+ years in distributed systems and AI infrastructure",
        "Former lead architect at Hyperscale Cloud Corp",
        "Author of 'Scalable Intelligence' (O'Reilly, 2025)",
        "Topic: Building Resilient AI Systems at Global Scale",
    ])

    # --- Slide 4: Industry Trends ---
    slide4 = prs.slides.add_slide(content_layout)
    add_title_content(slide4, "Industry Trends 2026", [
        "Edge computing adoption up 340% year-over-year",
        "AI agent frameworks becoming enterprise standard",
        "Zero-trust architecture now mandatory for compliance",
        "WebAssembly moving beyond browser to server-side",
        "Quantum-resistant cryptography integration accelerating",
        "Sustainable computing initiatives driving hardware design",
    ])

    # --- Slide 5: Technical Workshops ---
    slide5 = prs.slides.add_slide(content_layout)
    add_title_content(slide5, "Workshop Tracks", [
        "Track A: Cloud-Native Architecture Patterns",
        "  - Kubernetes operator design for stateful workloads",
        "  - Service mesh observability with OpenTelemetry",
        "  - Multi-cluster federation strategies",
        "Track B: ML Pipeline Engineering",
        "  - Feature store design and governance",
        "  - Model serving at 10M+ requests per second",
        "  - A/B testing infrastructure for ML models",
    ])

    # --- Slide 6: Speaker Lineup ---
    slide6 = prs.slides.add_slide(content_layout)
    add_title_content(slide6, "Featured Speakers", [
        "Raj Patel - VP Engineering, StreamForge",
        "Lena Vasquez - Principal Architect, DataWeave",
        "Chen Wei - Head of ML Platform, TechVault",
        "Samira Hadid - Director of Security, CyberShield",
        "Marcus Obi - CTO, GreenCompute Initiative",
        "Yuki Tanaka - Research Lead, QuantumBridge Labs",
    ])

    # --- Slide 7: Demo Showcase ---
    slide7 = prs.slides.add_slide(content_layout)
    add_title_content(slide7, "Live Demo Showcase", [
        "Real-time anomaly detection on streaming IoT data",
        "Autonomous code review agent with semantic understanding",
        "Cross-cloud data pipeline orchestration platform",
        "AR-assisted infrastructure monitoring dashboard",
        "Natural language database query interface",
    ])

    # --- Slide 8: Panel Discussion ---
    slide8 = prs.slides.add_slide(content_layout)
    add_title_content(slide8, "Panel: AI Ethics in Enterprise", [
        "Moderator: Prof. Elena Marchetti, Stanford AI Lab",
        "Key questions:",
        "  - How should companies audit AI decision-making?",
        "  - Balancing innovation speed with responsible deployment",
        "  - Regulatory landscape: EU AI Act implications",
        "  - Building inclusive AI teams and reducing bias",
    ])

    # --- Slide 9: Networking Events ---
    slide9 = prs.slides.add_slide(content_layout)
    add_title_content(slide9, "Networking Opportunities", [
        "Morning Coffee & Pastries (08:30 - 09:00)",
        "Speed Networking Session (12:30 - 13:00)",
        "Startup Alley - Meet 25 emerging tech companies",
        "Birds-of-a-Feather tables during lunch",
        "Evening Reception with live music (18:00 - 21:00)",
        "After-party at The Innovation Hub rooftop",
    ])

    # --- Slide 10: Sponsor Recognition ---
    slide10 = prs.slides.add_slide(content_layout)
    add_title_content(slide10, "Our Sponsors", [
        "Platinum: NexusAI Systems, CloudPeak Technologies",
        "Gold: DataWeave, StreamForge, QuantumBridge",
        "Silver: GreenCompute, CyberShield, CodeForge",
        "Community Partners: Open Source Alliance, DevMentor",
        "Media Partners: TechInsight Magazine, The Dev Report",
    ])

    # --- Slide 11: Venue Information ---
    slide11 = prs.slides.add_slide(content_layout)
    add_title_content(slide11, "Venue & Logistics", [
        "Grand Innovation Center, 450 Pacific Avenue, San Francisco",
        "Main Hall capacity: 2,500 attendees",
        "6 breakout rooms for workshops (50-150 capacity)",
        "Free parking in adjacent garage (code: TECH2026)",
        "BART: Montgomery Station, 5-minute walk",
        "Hotel block: Marriott Union Square ($189/night with code)",
    ])

    # --- Slide 12: Conference Stats ---
    slide12 = prs.slides.add_slide(content_layout)
    add_title_content(slide12, "By The Numbers", [
        "3,200+ registered attendees from 45 countries",
        "85 speakers across 4 tracks",
        "24 hands-on workshops",
        "50+ demo stations in the exhibition hall",
        "12 hours of recorded content available post-event",
        "$2.3M in total sponsor investment",
    ])

    # --- Slide 13: Call to Action ---
    slide13 = prs.slides.add_slide(content_layout)
    add_title_content(slide13, "Get Involved", [
        "Submit your lightning talk proposal by April 15",
        "Volunteer opportunities available - free admission",
        "Student discount: 60% off with valid .edu email",
        "Group rates: 5+ tickets at 25% discount",
        "Early bird pricing ends March 31, 2026",
        "Follow us: @TechConf2026 on all platforms",
    ])

    # --- Slide 14: Thank You ---
    slide14 = prs.slides.add_slide(content_layout)
    add_title_content(slide14, "Thank You!", [
        "We look forward to seeing you at Tech Conference 2026",
        "Questions? Email: info@techconf2026.io",
        "Website: www.techconference2026.com",
        "Registration: register.techconf2026.com",
        "#TechConf2026 #Innovation #AI #CloudNative",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
