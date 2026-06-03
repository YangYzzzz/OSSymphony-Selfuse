"""
Reward Script: Dark-themed presentation master
Task ID: impress_gf4_016
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Master background gradient #0F172A -> #1E293B
  Component 2 (0.25): Title placeholder text white (#FFFFFF) 40pt bold Calibri
  Component 3 (0.25): Content placeholder text #CBD5E1 18pt
  Component 4 (0.25): Blue horizontal line (#3B82F6) below title on all slides
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_016'

# Save any unsaved GUI state before verification
def persist_app_state():
    try:
        os.environ["DISPLAY"] = ":0"
        import pyautogui, time
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify dark-themed presentation master with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Pt, Emu
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: 14 slides
    num_slides = len(prs.slides)
    if num_slides != 14:
        print(f"WARN: Expected 14 slides, found {num_slides}")

    # =========================================================
    # Component 1: Master background gradient (0.25 points)
    # Task requires gradient from #0F172A (top) to #1E293B (bottom)
    # Initial has plain white background — this checks the change.
    # =========================================================
    try:
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'

        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slideMasters/slideMaster1.xml') as f:
                master_xml = f.read().decode()

        master_root = ET.fromstring(master_xml)
        bg = master_root.find(f'.//{{{ns_p}}}bg')

        gradient_ok = False
        if bg is not None:
            grad_fill = bg.find(f'.//{{{ns_a}}}gradFill')
            if grad_fill is not None:
                gs_list = grad_fill.findall(f'.//{{{ns_a}}}gs')
                stops = {}
                for gs in gs_list:
                    pos = gs.get('pos', '')
                    srgb = gs.find(f'{{{ns_a}}}srgbClr')
                    if srgb is not None:
                        stops[pos] = srgb.get('val', '').upper()

                # Check gradient colors (allow either order mapping)
                top_color = stops.get('0', '').upper()
                bottom_color = stops.get('100000', '').upper()

                if top_color == '0F172A' and bottom_color == '1E293B':
                    gradient_ok = True
                elif top_color == '1E293B' and bottom_color == '0F172A':
                    # Check if linear angle reverses direction
                    lin = grad_fill.find(f'{{{ns_a}}}lin')
                    if lin is not None:
                        ang = lin.get('ang', '')
                        # 16200000 = bottom-to-top (reversed)
                        if ang == '16200000':
                            gradient_ok = True

                print(f"  Gradient stops found: {stops}")

        if gradient_ok:
            print(f"PASS: Component 1 — Master gradient #0F172A -> #1E293B (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — Master gradient not matching #0F172A -> #1E293B")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # =========================================================
    # Component 2: Title text white (#FFFFFF) 40pt bold Calibri (0.25 points)
    # Initial has title in black, no explicit size/bold.
    # Check across multiple slides that title runs have the correct properties.
    # =========================================================
    try:
        title_pass_count = 0
        title_check_count = 0

        for si, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not hasattr(shape, 'text_frame'):
                    continue
                # Identify title placeholders
                ph = shape.placeholder_format
                if ph is not None and ph.idx == 0:  # idx 0 = title
                    for para in shape.text_frame.paragraphs:
                        runs = [r for r in para.runs if (r.text or "").strip()]
                        for run in runs:
                            title_check_count += 1
                            checks_passed = 0
                            total_checks = 4

                            # Check color is white
                            try:
                                if run.font.color.type is not None and str(run.font.color.rgb).upper() == 'FFFFFF':
                                    checks_passed += 1
                            except:
                                pass

                            # Check size is 40pt (508000 EMU)
                            if run.font.size is not None and abs(run.font.size - 508000) < 10000:
                                checks_passed += 1

                            # Check bold
                            if run.font.bold is True:
                                checks_passed += 1

                            # Check font is Calibri (or Montserrat)
                            font_name = (run.font.name or '').lower()
                            if font_name in ('calibri', 'montserrat'):
                                checks_passed += 1

                            if checks_passed >= 3:  # Allow one property to be inherited
                                title_pass_count += 1

        if title_check_count > 0 and title_pass_count >= min(title_check_count, num_slides * 0.5):
            print(f"PASS: Component 2 — Title: white, 40pt, bold, Calibri ({title_pass_count}/{title_check_count} runs pass) (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 — Title formatting: {title_pass_count}/{title_check_count} runs pass")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # =========================================================
    # Component 3: Content placeholder text #CBD5E1 18pt (0.25 points)
    # Initial has content in black, no explicit size.
    # Check body/subtitle placeholders across slides.
    # =========================================================
    try:
        content_pass_count = 0
        content_check_count = 0

        for si, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not hasattr(shape, 'text_frame'):
                    continue
                ph = shape.placeholder_format
                if ph is not None and ph.idx in (1, 10, 11, 12, 13, 14):  # body/content indices
                    for para in shape.text_frame.paragraphs:
                        runs = [r for r in para.runs if (r.text or "").strip()]
                        for run in runs:
                            content_check_count += 1
                            checks_passed = 0
                            total_checks = 2

                            # Check color is #CBD5E1
                            try:
                                if run.font.color.type is not None and str(run.font.color.rgb).upper() == 'CBD5E1':
                                    checks_passed += 1
                            except:
                                pass

                            # Check size is 18pt (228600 EMU)
                            if run.font.size is not None and abs(run.font.size - 228600) < 10000:
                                checks_passed += 1

                            if checks_passed >= 1:  # At least color OR size correct
                                content_pass_count += 1

        if content_check_count > 0 and content_pass_count >= min(content_check_count, num_slides * 0.5):
            print(f"PASS: Component 3 — Content: #CBD5E1, 18pt ({content_pass_count}/{content_check_count} runs pass) (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 — Content formatting: {content_pass_count}/{content_check_count} runs pass")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # =========================================================
    # Component 4: Blue horizontal line (#3B82F6) on slides (0.25 points)
    # Initial has no line shapes. Golden has cxnSp on every slide.
    # =========================================================
    try:
        slides_with_blue_line = 0

        with zipfile.ZipFile(file_path, 'r') as zf:
            for si in range(1, num_slides + 1):
                slide_xml_path = f'ppt/slides/slide{si}.xml'
                try:
                    with zf.open(slide_xml_path) as f:
                        slide_xml = f.read().decode()

                    slide_root = ET.fromstring(slide_xml)
                    spTree = slide_root.find(f'.//{{{ns_p}}}cSld/{{{ns_p}}}spTree')
                    if spTree is None:
                        continue

                    # Look for connector shapes (cxnSp) or regular shapes with line geometry
                    found_blue_line = False

                    # Check cxnSp elements
                    for cxn in spTree.findall(f'{{{ns_p}}}cxnSp'):
                        srgb = cxn.find(f'.//{{{ns_a}}}ln/{{{ns_a}}}solidFill/{{{ns_a}}}srgbClr')
                        if srgb is not None and srgb.get('val', '').upper() == '3B82F6':
                            found_blue_line = True
                            break

                    # Also check sp elements with line preset geometry
                    if not found_blue_line:
                        for sp in spTree.findall(f'{{{ns_p}}}sp'):
                            geom = sp.find(f'.//{{{ns_a}}}prstGeom')
                            if geom is not None and geom.get('prst') == 'line':
                                srgb = sp.find(f'.//{{{ns_a}}}ln/{{{ns_a}}}solidFill/{{{ns_a}}}srgbClr')
                                if srgb is not None and srgb.get('val', '').upper() == '3B82F6':
                                    found_blue_line = True
                                    break

                    # Also check outline/solidFill on shapes named with "line" in the name
                    if not found_blue_line:
                        for sp in spTree.findall(f'{{{ns_p}}}sp'):
                            srgb_fills = sp.findall(f'.//{{{ns_a}}}solidFill/{{{ns_a}}}srgbClr')
                            for srgb in srgb_fills:
                                if srgb.get('val', '').upper() == '3B82F6':
                                    # Confirm it's a line-like shape (height=0 or very small)
                                    ext = sp.find(f'.//{{{ns_a}}}ext')
                                    if ext is not None:
                                        cy = int(ext.get('cy', '999999'))
                                        if cy < 50000:  # very thin = line
                                            found_blue_line = True
                                            break

                    if found_blue_line:
                        slides_with_blue_line += 1

                except KeyError:
                    pass

        # Require blue line on at least 10 of 14 slides (allow some tolerance)
        threshold = max(1, int(num_slides * 0.7))
        if slides_with_blue_line >= threshold:
            print(f"PASS: Component 4 — Blue line on {slides_with_blue_line}/{num_slides} slides (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 4 — Blue line on {slides_with_blue_line}/{num_slides} slides (need >= {threshold})")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
