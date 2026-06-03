"""
FINAL REWARD SCRIPT - SUCCESS
Task: Insert a bookmark named 'sec-methods' at Heading 2 'Methods'.
Generated: 2025-10-17 08:01:29
Status: success
Model: azure-o3
Total Steps: 9
"""

###############################################################################
# Reward Script for: Insert a bookmark named 'sec-methods' at Heading 2 'Methods'
###############################################################################
# This script automatically verifies that the learner inserted a bookmark named
# "sec-methods" on the slide containing the Heading 2 text "Methods" in a PPTX
# presentation located in /home/user.
#
# Scoring (progressive):
#   0.4 – Slide with heading "Methods" (Heading 2) exists
#   0.6 – A shape on that same slide is named exactly "sec-methods"
#   1.0 – Both conditions satisfied
###############################################################################

import os
from pptx import Presentation


def locate_pptx(search_dir: str = "/home/user") -> str | None:
    """Locate the target PPTX file in the given directory.

    Preference rules:
      1. If only one *.pptx exists, use it.
      2. If multiple, prefer one that does NOT contain the word "golden".
      3. Otherwise, pick the first alphabetically.
    """
    pptx_files = [f for f in os.listdir(search_dir)
                  if f.lower().endswith('.pptx') and not f.startswith('~$')]
    if not pptx_files:
        return None
    if len(pptx_files) == 1:
        return os.path.join(search_dir, pptx_files[0])

    non_golden = [f for f in pptx_files if 'golden' not in f.lower()]
    chosen = sorted(non_golden)[0] if non_golden else sorted(pptx_files)[0]
    return os.path.join(search_dir, chosen)


def verify_bookmark_methods(pptx_path: str) -> float:
    """Verify task completion and return a score between 0.0 and 1.0."""

    max_score = 1.0
    score = 0.0

    # 1) Load presentation ----------------------------------------------------
    try:
        prs = Presentation(pptx_path)
    except Exception as exc:
        print(f"✗ Failed to open presentation: {exc}")
        print("REWARD: 0.0")
        return 0.0

    print(f"✓ Loaded presentation: {pptx_path} (slides: {len(prs.slides)})")

    # 2) Locate Heading 2 "Methods" -----------------------------------------
    methods_slide_idx = None  # slide index where heading is found
    heading_shape = None      # shape containing the heading text

    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if getattr(shape, 'has_text_frame', False):
                if shape.text.strip().lower() == 'methods':
                    methods_slide_idx = idx
                    heading_shape = shape
                    break
        if methods_slide_idx is not None:
            break

    if methods_slide_idx is not None:
        print(f"✓ Found heading 'Methods' on slide {methods_slide_idx + 1}")
        score += 0.4
    else:
        print("✗ Heading 'Methods' not found in any slide")

    # 3) Verify bookmark named "sec-methods" on same slide -------------------
    bookmark_found = False
    if methods_slide_idx is not None:
        slide = prs.slides[methods_slide_idx]
        for shape in slide.shapes:
            if shape.name and shape.name.strip().lower() == 'sec-methods':
                bookmark_found = True
                break

        if bookmark_found:
            print("✓ Bookmark named 'sec-methods' found on Methods slide")
            score += 0.6
        else:
            print("✗ Bookmark 'sec-methods' NOT found on Methods slide")

    final_score = min(max_score, round(score, 2))
    print(f"REWARD: {final_score}")
    return final_score


# ------------------------- MAIN EXECUTION BLOCK -----------------------------
if __name__ == "__main__":
    pptx_file = locate_pptx()
    if pptx_file is None:
        print("✗ No .pptx file found in /home/user")
        print("REWARD: 0.0")
    else:
        verify_bookmark_methods(pptx_file)

