"""
Initial Setup: Create a 5-slide Budget Overview presentation with slide 2 ready for chart creation.
Task ID: impress_gf2_022
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_022'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Annual Budget Overview"
    slide1.placeholders[1].text = "Fiscal Year 2024 — Executive Summary"

    # --- Slide 2: Department Budget Allocation (empty content area, no chart) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title text box at top
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Department Budget Allocation"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x3A, 0x5C)

    # --- Slide 3: Revenue Trends ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    tb3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Revenue Trends — Q1 to Q4"
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.runs[0]
    r3.font.name = "Arial"
    r3.font.size = Pt(28)
    r3.font.bold = True
    r3.font.color.rgb = RGBColor(0x2E, 0x3A, 0x5C)

    # Add content text
    content3 = slide3.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5))
    ctf3 = content3.text_frame
    ctf3.word_wrap = True
    lines = [
        "Q1 Revenue: $2.4M (+8% YoY)",
        "Q2 Revenue: $2.7M (+12% YoY)",
        "Q3 Revenue: $3.1M (+15% YoY)",
        "Q4 Revenue (projected): $3.5M (+18% YoY)",
        "",
        "Key drivers: expansion into APAC markets, new enterprise contracts,",
        "and improved customer retention rates across all segments.",
    ]
    for i, line in enumerate(lines):
        if i == 0:
            ctf3.paragraphs[0].text = line
        else:
            p = ctf3.add_paragraph()
            p.text = line
        para = ctf3.paragraphs[i]
        for run in para.runs:
            run.font.name = "Arial"
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 4: Headcount Planning ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    tb4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf4 = tb4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Headcount Planning by Department"
    p4.alignment = PP_ALIGN.CENTER
    r4 = p4.runs[0]
    r4.font.name = "Arial"
    r4.font.size = Pt(28)
    r4.font.bold = True
    r4.font.color.rgb = RGBColor(0x2E, 0x3A, 0x5C)

    # Add a table with headcount data
    rows, cols = 6, 4
    table_shape = slide4.shapes.add_table(rows, cols, Inches(1.5), Inches(1.5), Inches(10), Inches(4))
    table = table_shape.table
    headers = ["Department", "Current FTE", "Planned FTE", "Change"]
    dept_data = [
        ["HR", "42", "48", "+6"],
        ["IT", "105", "120", "+15"],
        ["Marketing", "68", "75", "+7"],
        ["Operations", "130", "140", "+10"],
        ["R&D", "55", "62", "+7"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row_data in enumerate(dept_data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 5: Next Steps ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    tb5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf5 = tb5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Next Steps & Action Items"
    p5.alignment = PP_ALIGN.CENTER
    r5 = p5.runs[0]
    r5.font.name = "Arial"
    r5.font.size = Pt(28)
    r5.font.bold = True
    r5.font.color.rgb = RGBColor(0x2E, 0x3A, 0x5C)

    content5 = slide5.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5))
    ctf5 = content5.text_frame
    ctf5.word_wrap = True
    items = [
        "1. Finalize department budget allocations by March 15, 2024",
        "2. Submit headcount increase requests to Finance by March 22",
        "3. Complete vendor contract reviews for Q2 procurement",
        "4. Schedule departmental budget review meetings (April 1-5)",
        "5. Present final budget proposal to Board of Directors on April 12",
    ]
    for i, item in enumerate(items):
        if i == 0:
            ctf5.paragraphs[0].text = item
        else:
            p = ctf5.add_paragraph()
            p.text = item
        para = ctf5.paragraphs[i]
        for run in para.runs:
            run.font.name = "Arial"
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
