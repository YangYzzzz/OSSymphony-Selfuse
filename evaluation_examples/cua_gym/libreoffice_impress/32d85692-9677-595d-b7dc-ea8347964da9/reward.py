"""
Reward Script: Create embedded donut chart on slide 2 with budget allocation
Task ID: impress_gf2_022
Domain: libreoffice_impress
Scoring:
  Component 1: Donut chart exists on slide 2          — 0.20 pts
  Component 2: Correct categories and values           — 0.25 pts
  Component 3: Data labels show percentages (XML)      — 0.20 pts
  Component 4: 5 distinct colors on segments           — 0.10 pts
  Component 5: Legend is present                       — 0.10 pts
  Component 6: 'Budget 2024' text box on slide 2       — 0.15 pts
  Total: 1.00
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_022'


def persist_app_state(domain: str):
    """Best-effort save of any open LibreOffice document."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.chart import XL_CHART_TYPE
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 2")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]  # 0-indexed, slide 2

    # Find chart shape on slide 2
    chart_shape = None
    for shape in slide2.shapes:
        if shape.has_chart:
            chart_shape = shape
            break

    # ============================================================
    # Component 1: Donut chart exists on slide 2 (0.20 points)
    # ============================================================
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            # Check it is a DOUGHNUT type (-4120)
            if chart.chart_type == XL_CHART_TYPE.DOUGHNUT:
                print(f"PASS: Component 1 — Donut chart found on slide 2 (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 1 — Chart found but type is {chart.chart_type}, expected DOUGHNUT")
        else:
            print(f"FAIL: Component 1 — No chart found on slide 2")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # ============================================================
    # Component 2: Correct categories and values (0.25 points)
    # Expected: HR=15, IT=25, Marketing=20, Operations=30, R&D=10
    # ============================================================
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            plot = chart.plots[0]
            categories = [str(c) for c in plot.categories]
            values = list(chart.series[0].values)

            expected_data = {
                'HR': 15.0,
                'IT': 25.0,
                'Marketing': 20.0,
                'Operations': 30.0,
                'R&D': 10.0,
            }

            # Check categories match (order-independent)
            cats_match = set(categories) == set(expected_data.keys())
            # Check values match for each category
            vals_match = False
            if cats_match and len(categories) == len(values):
                cat_val_map = dict(zip(categories, values))
                vals_match = all(
                    abs(cat_val_map.get(k, -1) - v) < 0.01
                    for k, v in expected_data.items()
                )

            if cats_match and vals_match:
                print(f"PASS: Component 2 — Categories {categories} and values {values} match expected (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 2 — Categories: {categories}, Values: {values}. Expected: {expected_data}")
        else:
            print(f"FAIL: Component 2 — No chart to check categories/values")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # ============================================================
    # Component 3: Data labels show percentages (0.20 points)
    # Check via XML for showPercent attribute
    # ============================================================
    try:
        found_show_percent = False
        with zipfile.ZipFile(file_path, 'r') as zf:
            chart_files = [n for n in zf.namelist() if 'chart' in n.lower() and n.endswith('.xml')]
            for cf in chart_files:
                with zf.open(cf) as f:
                    root = ET.fromstring(f.read())
                    # Look for showPercent elements with val="1" or val="true"
                    for elem in root.iter():
                        tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
                        if tag == 'showPercent':
                            val = elem.get('val', '0')
                            if val in ('1', 'true'):
                                found_show_percent = True
                                break
                if found_show_percent:
                    break

        if found_show_percent:
            print(f"PASS: Component 3 — Data labels show percentages (showPercent=1) (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 3 — showPercent not enabled in chart XML")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # ============================================================
    # Component 4: 5 distinct colors on chart segments (0.10 points)
    # Check via XML for dPt fill colors
    # ============================================================
    try:
        colors = set()
        with zipfile.ZipFile(file_path, 'r') as zf:
            chart_files = [n for n in zf.namelist() if 'chart' in n.lower() and n.endswith('.xml')]
            for cf in chart_files:
                with zf.open(cf) as f:
                    root = ET.fromstring(f.read())
                    # Find all dPt elements and extract srgbClr from their fill
                    in_dpt = False
                    for elem in root.iter():
                        tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
                        if tag == 'dPt':
                            in_dpt = True
                        if tag == 'srgbClr' and in_dpt:
                            colors.add(elem.get('val', '').upper())

        # We need at least 5 distinct colors
        if len(colors) >= 5:
            print(f"PASS: Component 4 — {len(colors)} distinct colors found: {colors} (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 4 — Only {len(colors)} distinct colors found: {colors}, need 5")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # ============================================================
    # Component 5: Legend is present on chart (0.10 points)
    # ============================================================
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            if chart.has_legend:
                print(f"PASS: Component 5 — Chart has legend (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 5 — Chart has no legend")
        else:
            print(f"FAIL: Component 5 — No chart to check legend")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # ============================================================
    # Component 6: 'Budget 2024' text box on slide 2 (0.15 points)
    # This is the center text label in the donut hole
    # ============================================================
    try:
        found_budget_text = False
        for shape in slide2.shapes:
            if shape.has_text_frame:
                full_text = ''.join(p.text for p in shape.text_frame.paragraphs).strip()
                if 'Budget 2024' in full_text:
                    found_budget_text = True
                    break

        if found_budget_text:
            print(f"PASS: Component 6 — 'Budget 2024' text box found on slide 2 (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 6 — 'Budget 2024' text not found on slide 2")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice
persist_app_state("libreoffice_impress")

# Execute verification
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
