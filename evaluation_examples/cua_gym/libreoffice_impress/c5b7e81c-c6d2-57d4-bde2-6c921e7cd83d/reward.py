"""
FINAL REWARD SCRIPT - SUCCESS
Task: I've placed three images on my slide, and I'd love to see them spaced evenly across the width. How do I set them up so they're distributed equally horizontally?
Generated: 2025-08-07 09:19:25
Status: success
Model: o4-mini
Total Steps: 1
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def verify_distribution(file_path):
    print("Checking task: images distributed evenly horizontally")
    total_score = 0.0
    max_score = 1.0

    # Requirement 1: File exists
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print(f"REWARD: {total_score}")
        return total_score
    print("✓ File exists (0.2 points)")
    total_score += 0.2

    # Load presentation
    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded successfully with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        print(f"REWARD: {total_score}")
        return total_score

    # Consider first slide for image distribution
    slide = prs.slides[0]

    # Requirement 2: Exactly three pictures on slide
    pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    pic_count = len(pictures)
    print(f"Found {pic_count} picture(s) on the slide")
    if pic_count == 3:
        print("✓ Correct number of images (3) (0.4 points)")
        total_score += 0.4
    elif pic_count > 3:
        print("✓ More than three images found, using first three shapes for spacing check (0.2 points)")
        total_score += 0.2
    else:
        print("✗ Incorrect number of images (<3)")
        print(f"REWARD: {total_score}")
        return total_score

    # Sorting first three images by left position
    pics_sorted = sorted(pictures[:3], key=lambda p: p.left)
    lefts = [p.left for p in pics_sorted]
    rights = [p.left + p.width for p in pics_sorted]

    # Compute gaps between image edges
    gap1 = lefts[1] - rights[0]
    gap2 = lefts[2] - rights[1]
    print(f"Gap between image1 and image2: {gap1}")
    print(f"Gap between image2 and image3: {gap2}")

    # Requirement 3: Gaps equal within tolerance
    tolerance = 50000  # EMU tolerance (~0.05 inches)
    if abs(gap1 - gap2) <= tolerance:
        print(f"✓ Gaps are equal within tolerance ({tolerance}) (0.4 points)")
        total_score += 0.4
    else:
        print(f"✗ Gaps are not equal (diff={abs(gap1-gap2)} > {tolerance}) (0 points)")
        # No partial gap points

    # Final score
    final_score = min(total_score, max_score)
    print(f"Final calculated score: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score

# Execute verification
if __name__ == '__main__':
    file_path = '/home/user/ive_placed_three_images_on_my_slide_and_id_love_to_see_them_spaced_evenly_across_the_width_how_do_i_.pptx'
    verify_distribution(file_path)
