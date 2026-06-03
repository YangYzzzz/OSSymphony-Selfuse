"""
Reward Script: Clear all notes from the entire presentation
Task ID: impress_ndo_035
Domain: libreoffice_impress
Scoring:
  Precondition gate: File exists, loads, has 20 slides with content (no points)
  Component 1 (0.5): All 20 slides have empty notes
  Component 2 (0.5): Progressive credit for each of the 15 originally-noted slides cleared
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_035'


def persist_app_state(domain: str):
    """Save any unsaved GUI state before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_slide_notes(slide):
    """Get notes text for a slide, returning empty string if no notes."""
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Precondition gate: File must have 20 slides with content
    # This is NOT scored — it's a precondition that should be true in both envs
    if num_slides != 20:
        print(f"PRECONDITION FAIL: Expected 20 slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    slides_with_content = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                slides_with_content += 1
                break

    if slides_with_content < 15:
        print(f"PRECONDITION FAIL: Only {slides_with_content} slides have content (expected >= 15)")
        print("REWARD: 0.0")
        return 0.0

    print(f"PRECONDITION OK: {num_slides} slides, {slides_with_content} with content")

    # Component 1: All 20 slides have empty notes (0.5 points)
    # This is the core task requirement: "Every slide should have empty notes when done"
    # Only awards if ALL slides are clear — no partial credit here
    try:
        slides_with_notes = []
        for i, slide in enumerate(prs.slides):
            notes_text = get_slide_notes(slide)
            if notes_text:
                slides_with_notes.append(i + 1)

        if len(slides_with_notes) == 0:
            print(f"PASS: Component 1 — All {num_slides} slides have empty notes (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 1 — {len(slides_with_notes)} slides still have notes: {slides_with_notes}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Progressive credit for clearing the 15 originally-noted slides (0.5 points)
    # Slides 1-15 originally had notes. Each cleared slide earns proportional credit.
    # This component provides partial credit for incomplete task completion.
    try:
        originally_with_notes = list(range(0, 15))  # 0-indexed: slides 1-15
        cleared_count = 0
        for idx in originally_with_notes:
            slide = prs.slides[idx]
            notes_text = get_slide_notes(slide)
            if not notes_text:
                cleared_count += 1
            else:
                print(f"  Slide {idx+1} still has notes: {notes_text[:60]}...")

        if cleared_count == 15:
            print(f"PASS: Component 2 — All 15 originally-noted slides cleared (0.5 pts)")
            total_score += 0.5
        elif cleared_count > 0:
            partial = round(0.5 * (cleared_count / 15), 4)
            print(f"PARTIAL: Component 2 — {cleared_count}/15 originally-noted slides cleared ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — 0/15 originally-noted slides cleared")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist app state if needed, then verify
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
