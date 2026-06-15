"""
Initial Setup: Clear all notes from the entire presentation
Task ID: impress_ndo_035
Domain: libreoffice_impress

Creates a 20-slide business presentation where 15 slides have speaker notes.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_035'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=None, alignment=None):
    """Helper to add a textbox with formatted text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q4 2025 Strategic Review"
    slide1.placeholders[1].text = "Nextera Solutions Inc."
    slide1.notes_slide.notes_text_frame.text = "Welcome everyone to the Q4 strategic review. Please hold questions until the end of each section."

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "1. Financial Highlights"
    tf2.add_paragraph().text = "2. Product Roadmap Update"
    tf2.add_paragraph().text = "3. Market Expansion Plans"
    tf2.add_paragraph().text = "4. Team & Hiring"
    tf2.add_paragraph().text = "5. Q1 2026 Objectives"
    slide2.notes_slide.notes_text_frame.text = "Estimated time: 45 minutes total. Financial section ~15 min, product ~10 min, market ~10 min, team ~5 min, objectives ~5 min."

    # --- Slide 3: Financial Highlights ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Financial Highlights"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Revenue: $12.4M (+18% YoY)"
    tf3.add_paragraph().text = "Gross Margin: 72.3%"
    tf3.add_paragraph().text = "EBITDA: $3.1M"
    tf3.add_paragraph().text = "Cash Position: $8.7M"
    tf3.add_paragraph().text = "Customer Count: 347 (+42 net new)"
    slide3.notes_slide.notes_text_frame.text = "Highlight that revenue growth accelerated from 14% in Q3. Gross margin improvement driven by infrastructure optimization completed in September."

    # --- Slide 4: Revenue Breakdown ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Revenue Breakdown by Segment"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Enterprise: $7.2M (58%)"
    tf4.add_paragraph().text = "Mid-Market: $3.8M (31%)"
    tf4.add_paragraph().text = "SMB: $1.4M (11%)"
    slide4.notes_slide.notes_text_frame.text = "Enterprise segment grew 24% due to the Meridian Corp and Astral Health deals. Mid-market is our fastest-growing segment by percentage at 31% YoY."

    # --- Slide 5: Customer Retention ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Customer Retention Metrics"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Net Revenue Retention: 118%"
    tf5.add_paragraph().text = "Logo Retention: 94%"
    tf5.add_paragraph().text = "Average Contract Value: $35.7K"
    tf5.add_paragraph().text = "Expansion Revenue: $1.8M"
    slide5.notes_slide.notes_text_frame.text = "NRR above 115% is best-in-class for B2B SaaS. The three churned accounts were all SMB tier and cited budget constraints rather than product issues."

    # --- Slide 6: Product Roadmap ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Product Roadmap - Delivered in Q4"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "AI-Powered Analytics Dashboard (shipped Oct 15)"
    tf6.add_paragraph().text = "SSO Integration with Okta & Azure AD (shipped Nov 2)"
    tf6.add_paragraph().text = "Mobile App v2.0 (shipped Dec 8)"
    tf6.add_paragraph().text = "API Rate Limit Improvements (shipped Nov 20)"
    slide6.notes_slide.notes_text_frame.text = "The AI analytics dashboard has been the most requested feature since Q2. Early adoption rate is 67% among enterprise customers within the first month."

    # --- Slide 7: Upcoming Features ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Product Roadmap - Q1 2026"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "Custom Report Builder (target: Jan 30)"
    tf7.add_paragraph().text = "Webhook Event System (target: Feb 15)"
    tf7.add_paragraph().text = "Multi-tenant Administration (target: Mar 10)"
    tf7.add_paragraph().text = "Data Export API v3 (target: Mar 28)"
    slide7.notes_slide.notes_text_frame.text = "Custom report builder is the top-priority item. Meridian Corp has made this a condition for their renewal in March. Engineering team is confident in the January timeline."

    # --- Slide 8: Technical Debt ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Technical Infrastructure"
    tf8 = slide8.placeholders[1].text_frame
    tf8.text = "Cloud costs reduced 22% via reserved instances"
    tf8.add_paragraph().text = "Average API latency: 45ms (from 78ms)"
    tf8.add_paragraph().text = "99.97% uptime in Q4"
    tf8.add_paragraph().text = "Security audit completed — no critical findings"
    slide8.notes_slide.notes_text_frame.text = "The latency improvement came from the database sharding project that Raj's team completed ahead of schedule. Uptime was affected by one 12-minute outage in November caused by a DNS provider issue."

    # --- Slide 9: Market Expansion ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Market Expansion - EMEA Launch"
    tf9 = slide9.placeholders[1].text_frame
    tf9.text = "London office opened October 1"
    tf9.add_paragraph().text = "8 enterprise prospects in pipeline"
    tf9.add_paragraph().text = "First EMEA customer: Thornfield Industries (signed Nov)"
    tf9.add_paragraph().text = "GDPR compliance audit passed"
    slide9.notes_slide.notes_text_frame.text = "EMEA pipeline value is approximately $2.1M. Thornfield Industries is a strong reference account. We should mention the partnership with Denholm Consulting for local support."

    # --- Slide 10: Competitive Landscape ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "Competitive Landscape"
    tf10 = slide10.placeholders[1].text_frame
    tf10.text = "Win rate vs. Compex Analytics: 62% (up from 55%)"
    tf10.add_paragraph().text = "Win rate vs. DataBridge: 71%"
    tf10.add_paragraph().text = "Key differentiators: AI features, implementation speed"
    tf10.add_paragraph().text = "Compex raised $40M Series C — expect feature parity push"
    slide10.notes_slide.notes_text_frame.text = "Compex's fundraise is concerning. They will likely invest heavily in AI features to match ours. We need to maintain our 6-month lead. DataBridge seems to be pivoting toward a different market segment."

    # --- Slide 11: Partner Ecosystem ---
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    slide11.shapes.title.text = "Partner Ecosystem"
    tf11 = slide11.placeholders[1].text_frame
    tf11.text = "12 new integration partners in Q4"
    tf11.add_paragraph().text = "Salesforce connector GA (200+ installs)"
    tf11.add_paragraph().text = "HubSpot integration in beta"
    tf11.add_paragraph().text = "Partner-sourced revenue: $1.2M (10% of total)"
    slide11.notes_slide.notes_text_frame.text = "The Salesforce connector has been a significant driver of enterprise deals. Partner-sourced revenue target for 2026 is 20% of total revenue."

    # --- Slide 12: Team Overview ---
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])
    slide12.shapes.title.text = "Team & Organization"
    tf12 = slide12.placeholders[1].text_frame
    tf12.text = "Headcount: 89 (from 71 at start of Q4)"
    tf12.add_paragraph().text = "Engineering: 42 | Sales: 18 | CS: 14 | G&A: 15"
    tf12.add_paragraph().text = "Q4 Hires: 18 new team members"
    tf12.add_paragraph().text = "Voluntary attrition: 4.2% (annualized)"
    slide12.notes_slide.notes_text_frame.text = "Low attrition is a strong signal. Key hires include VP of Engineering (Priya Sharma from Stripe) and Director of EMEA Sales (James Whitfield from Tableau)."

    # --- Slide 13: Hiring Plan ---
    slide13 = prs.slides.add_slide(prs.slide_layouts[1])
    slide13.shapes.title.text = "Q1 2026 Hiring Plan"
    tf13 = slide13.placeholders[1].text_frame
    tf13.text = "10 Engineering roles (4 senior, 6 mid-level)"
    tf13.add_paragraph().text = "3 Sales Development Representatives"
    tf13.add_paragraph().text = "2 Customer Success Managers (EMEA)"
    tf13.add_paragraph().text = "1 Head of Product Marketing"
    slide13.notes_slide.notes_text_frame.text = "Budget approved for all 16 positions. Recruiting pipeline is strong for engineering roles. Sales hiring may take longer due to competition for talent in the London market."

    # --- Slide 14: Customer Testimonials ---
    slide14 = prs.slides.add_slide(prs.slide_layouts[1])
    slide14.shapes.title.text = "Customer Feedback"
    tf14 = slide14.placeholders[1].text_frame
    tf14.text = '"The AI analytics feature has transformed how we make decisions." — VP Data, Meridian Corp'
    tf14.add_paragraph().text = '"Implementation took 3 weeks instead of the usual 3 months." — CTO, Brightpath Health'
    tf14.add_paragraph().text = "NPS Score: 72 (up from 65 in Q3)"
    slide14.notes_slide.notes_text_frame.text = "Consider featuring the Meridian Corp quote in our next marketing campaign. Brightpath Health is open to being a case study."

    # --- Slide 15: Risk Assessment ---
    slide15 = prs.slides.add_slide(prs.slide_layouts[1])
    slide15.shapes.title.text = "Risk Assessment"
    tf15 = slide15.placeholders[1].text_frame
    tf15.text = "Competitive pressure from Compex (HIGH)"
    tf15.add_paragraph().text = "Key person dependency on 3 senior engineers (MEDIUM)"
    tf15.add_paragraph().text = "EMEA regulatory changes pending (MEDIUM)"
    tf15.add_paragraph().text = "Cloud cost inflation risk (LOW)"
    slide15.notes_slide.notes_text_frame.text = "The Compex risk is our top concern. Mitigation: accelerate AI roadmap and strengthen customer relationships. Key person risk is being addressed through knowledge-sharing initiatives and documentation sprints."

    # --- Slide 16: Q1 Objectives (no notes) ---
    slide16 = prs.slides.add_slide(prs.slide_layouts[1])
    slide16.shapes.title.text = "Q1 2026 Objectives"
    tf16 = slide16.placeholders[1].text_frame
    tf16.text = "Achieve $14M quarterly revenue target"
    tf16.add_paragraph().text = "Launch custom report builder"
    tf16.add_paragraph().text = "Close 5 EMEA enterprise accounts"
    tf16.add_paragraph().text = "Reduce churn to below 5%"

    # --- Slide 17: Key Metrics Dashboard (no notes) ---
    slide17 = prs.slides.add_slide(prs.slide_layouts[1])
    slide17.shapes.title.text = "Key Performance Indicators"
    tf17 = slide17.placeholders[1].text_frame
    tf17.text = "MRR Target: $4.7M"
    tf17.add_paragraph().text = "CAC Payback: < 12 months"
    tf17.add_paragraph().text = "LTV/CAC Ratio: > 5x"
    tf17.add_paragraph().text = "Support Ticket Resolution: < 4 hours"

    # --- Slide 18: Timeline (no notes) ---
    slide18 = prs.slides.add_slide(prs.slide_layouts[1])
    slide18.shapes.title.text = "Q1 2026 Timeline"
    tf18 = slide18.placeholders[1].text_frame
    tf18.text = "January: Custom Report Builder release"
    tf18.add_paragraph().text = "February: EMEA roadshow (London, Frankfurt, Paris)"
    tf18.add_paragraph().text = "March: Annual customer conference planning"
    tf18.add_paragraph().text = "March 31: Q1 board review"

    # --- Slide 19: Budget Summary (no notes) ---
    slide19 = prs.slides.add_slide(prs.slide_layouts[1])
    slide19.shapes.title.text = "Q1 2026 Budget Summary"
    tf19 = slide19.placeholders[1].text_frame
    tf19.text = "Total OpEx: $9.2M"
    tf19.add_paragraph().text = "Engineering: $4.1M | Sales & Marketing: $2.8M"
    tf19.add_paragraph().text = "G&A: $1.5M | EMEA Operations: $0.8M"
    tf19.add_paragraph().text = "Capital Reserve: $3.5M"

    # --- Slide 20: Thank You (no notes) ---
    slide20 = prs.slides.add_slide(prs.slide_layouts[0])
    slide20.shapes.title.text = "Thank You"
    slide20.placeholders[1].text = "Questions & Discussion"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
