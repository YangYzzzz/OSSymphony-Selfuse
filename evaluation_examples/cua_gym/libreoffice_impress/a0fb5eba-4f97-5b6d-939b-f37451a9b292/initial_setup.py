"""
Initial Setup: Create a branded presentation with 10 slides using 3 different master slide designs.
Task ID: impress_el_045
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_el_045'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Add a textbox with styled text to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def set_slide_background(slide, r, g, b):
    """Set solid background color for a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # We'll create slides with 3 distinct visual "master" styles applied manually
    # Style A: Dark blue background, white text (Corporate)
    # Style B: Light green background, dark text (Nature)
    # Style C: White background with red accent, dark text (Clean)

    slide_data = [
        # (slide_num, style, title, content)
        (1, 'A', 'Q4 2025 Strategic Review',
         'Acme Corporation | Board of Directors Presentation'),
        (2, 'A', 'Executive Summary',
         'Revenue grew 18% YoY reaching $142M in Q4.\nOperating margin improved to 23.5%.\nCustomer retention rate at 94.2%.'),
        (3, 'A', 'Financial Highlights',
         'Net Revenue: $142.3M (+18% YoY)\nGross Profit: $89.7M (63% margin)\nEBITDA: $41.2M (+22% YoY)\nFree Cash Flow: $28.5M'),
        (4, 'B', 'Sustainability Initiatives',
         'Carbon neutral operations by 2027.\nSolar panel deployment across 12 facilities.\nWaste reduction program saved 340 tons in Q4.\nEmployee green commute program: 67% participation.'),
        (5, 'B', 'Environmental Impact Report',
         'Water usage reduced by 25% vs Q3.\nRecycled materials in packaging: 82%.\nBiodiversity projects funded: 14 across 6 countries.\nGreen energy sourcing: 71% of total consumption.'),
        (6, 'B', 'Community Engagement',
         'STEM education sponsorship: 45 schools.\nLocal food bank partnership: 12,000 meals donated.\nVolunteer hours by employees: 8,400 in Q4.\nScholarship fund disbursed $320K to 40 students.'),
        (7, 'C', 'Product Roadmap 2026',
         'Platform v3.0 launch: March 2026.\nMobile app redesign: May 2026.\nAI-powered analytics module: July 2026.\nEnterprise SSO integration: September 2026.'),
        (8, 'C', 'Engineering Milestones',
         'API response time reduced to 45ms (from 120ms).\nInfrastructure migration to Kubernetes complete.\n99.97% uptime achieved in Q4.\nAutomated test coverage: 91%.'),
        (9, 'A', 'Market Expansion',
         'New offices in Berlin, Tokyo, and Sao Paulo.\nPartnership with Meridian Technologies signed.\nEnterprise clients added: 38 in Q4.\nTotal addressable market expanded to $4.2B.'),
        (10, 'C', 'Next Steps & Action Items',
         'Finalize 2026 budget allocation by Jan 15.\nLaunch customer advisory board in February.\nComplete Series D fundraising by Q1.\nHire 120 new engineers across 3 offices.'),
    ]

    for idx, (num, style, title, content) in enumerate(slide_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only layout

        # Apply background based on style
        if style == 'A':
            set_slide_background(slide, 0x1B, 0x2A, 0x4A)  # Dark navy blue
            title_color = RGBColor(0xFF, 0xFF, 0xFF)
            text_color = RGBColor(0xD0, 0xD8, 0xE8)
            accent_color = RGBColor(0x4E, 0xC9, 0xB0)
        elif style == 'B':
            set_slide_background(slide, 0xE8, 0xF5, 0xE9)  # Light green
            title_color = RGBColor(0x1B, 0x5E, 0x20)
            text_color = RGBColor(0x33, 0x33, 0x33)
            accent_color = RGBColor(0x2E, 0x7D, 0x32)
        else:  # C
            set_slide_background(slide, 0xFF, 0xFF, 0xFF)  # White
            title_color = RGBColor(0xC6, 0x28, 0x28)
            text_color = RGBColor(0x21, 0x21, 0x21)
            accent_color = RGBColor(0xC6, 0x28, 0x28)

        # Title
        if slide.shapes.title:
            slide.shapes.title.text = title
            for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                run.font.color.rgb = title_color
                run.font.size = Pt(32)
                run.font.bold = True

        # Accent line (top bar)
        from pptx.enum.shapes import MSO_SHAPE
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), prs.slide_width, Inches(0.08)
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = accent_color
        line_shape.line.fill.background()

        # Content text
        content_lines = content.split('\n')
        y_start = Inches(2.0)
        for i, line in enumerate(content_lines):
            add_textbox(
                slide,
                Inches(1.0), y_start + Inches(i * 0.65),
                Inches(11), Inches(0.6),
                line, font_size=16, color=text_color
            )

        # Slide number indicator
        add_textbox(
            slide,
            Inches(12.0), Inches(6.8),
            Inches(1), Inches(0.5),
            str(num), font_size=12, color=text_color,
            alignment=PP_ALIGN.RIGHT
        )

        # Style label in bottom-left (helps identify which "master" design)
        style_labels = {'A': 'Corporate', 'B': 'Nature', 'C': 'Clean'}
        add_textbox(
            slide,
            Inches(0.5), Inches(6.8),
            Inches(2), Inches(0.5),
            f'{style_labels[style]} Theme', font_size=10, color=text_color
        )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=3.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
