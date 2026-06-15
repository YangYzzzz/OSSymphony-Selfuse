"""
FINAL REWARD SCRIPT - SUCCESS
Task: Start Arabic page numbering at page 3 with value 1 and place it bottom-center.
Generated: 2025-10-17 07:19:45
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
import re
import math
from pptx import Presentation


def get_numeric_shapes(slide):
    """Return a list of (shape, integer_value) for shapes that contain ONLY an integer."""
    numeric = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if re.fullmatch(r"\d+", txt):
            numeric.append((shape, int(txt)))
    return numeric


def verify_arabic_page_numbering(file_path: str) -> float:
    """Verify that Arabic page numbering starts on slide 3 with 1 and is bottom-centered.

    Scoring (progressive):
        • 0.2 – First two slides have NO numeric page numbers.
        • 0.8 – Every subsequent slide i has number (i-2) placed bottom-center (within tolerance).
    Returns a float in [0,1]. Prints detailed diagnostics and the final REWARD value.
    """
    score = 0.0
    max_score = 1.0

    if not os.path.exists(file_path):
        print("✗ File not found:", file_path)
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print("✗ Unable to load PPTX:", e)
        print("REWARD: 0.0")
        return 0.0

    slide_w, slide_h = prs.slide_width, prs.slide_height
    total = len(prs.slides)
    print(f"✓ Loaded presentation ({total} slides)")

    # --- Requirement 1: Slides 1-2 must NOT have page numbers ---
    first_two_ok = True
    for idx in range(min(2, total)):
        nums = get_numeric_shapes(prs.slides[idx])
        if nums:
            first_two_ok = False
            vals = [v for _, v in nums]
            print(f"✗ Slide {idx+1} has unexpected numeric shapes: {vals}")
        else:
            print(f"✓ Slide {idx+1} has no page number (as expected)")

    if first_two_ok:
        score += 0.2

    # --- Requirement 2: Slides 3+ numbered sequentially starting with 1, bottom-center ---
    sequential_and_position_ok = True
    for s_idx in range(2, total):  # 0-based index, slide 3 == idx 2
        expected_val = s_idx - 1  # slide 3 -> 1, slide 4 -> 2, ...
        slide = prs.slides[s_idx]
        num_shapes = get_numeric_shapes(slide)

        if not num_shapes:
            sequential_and_position_ok = False
            print(f"✗ Slide {s_idx+1} missing page number")
            continue

        # pick numeric shape closest to bottom-center
        best_shape, best_val, best_dist = None, None, None
        for shp, val in num_shapes:
            cx = shp.left + shp.width / 2
            by = shp.top + shp.height  # bottom y
            dx = abs(cx - slide_w / 2) / slide_w
            dy = abs(by - slide_h) / slide_h
            dist = math.hypot(dx, dy)
            if best_dist is None or dist < best_dist:
                best_shape, best_val, best_dist = shp, val, dist

        # value check
        if best_val != expected_val:
            sequential_and_position_ok = False
            print(f"✗ Slide {s_idx+1} has number {best_val}, expected {expected_val}")
        else:
            print(f"✓ Slide {s_idx+1} correct number {best_val}")

        # position check (within 5% horizontally from center, bottom 10% vertically)
        cx_norm = (best_shape.left + best_shape.width / 2) / slide_w
        by_norm = (best_shape.top + best_shape.height) / slide_h
        center_ok = abs(cx_norm - 0.5) <= 0.05
        bottom_ok = by_norm >= 0.9
        if center_ok and bottom_ok:
            print(f"✓ Slide {s_idx+1} number positioned bottom-center")
        else:
            sequential_and_position_ok = False
            print(
                f"✗ Slide {s_idx+1} number mis-positioned: center_x={cx_norm:.2f}, bottom_y={by_norm:.2f}")

    if sequential_and_position_ok:
        score += 0.8

    final_score = min(score, max_score)
    print(f"Total Score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score


if __name__ == "__main__":
    # Path to the presentation created by the task performer
    FILE_PATH = "/home/user/start_arabic_page_numbering_at_page_3_with_value_1_and_place_it_bottom_center.pptx"
    verify_arabic_page_numbering(FILE_PATH)
