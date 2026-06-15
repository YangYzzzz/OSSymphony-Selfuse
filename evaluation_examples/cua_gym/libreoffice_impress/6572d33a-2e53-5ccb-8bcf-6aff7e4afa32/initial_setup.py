"""
Initial Setup: Create a 6-slide Growth Strategy presentation with an arrow on slide 4.
Task ID: impress_rp_015
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_015'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Growth Strategy 2026"
    slide1.placeholders[1].text = "Accelerating Market Expansion & Revenue Growth"

    # ---- Slide 2: Market Overview ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Market Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Global addressable market expanded to $4.2B in 2025"
    items2 = [
        "North America accounts for 38% of total revenue",
        "APAC region growing at 22% CAGR, fastest among all segments",
        "European regulatory changes opening new verticals",
        "Competitive landscape consolidating — top 5 players hold 61% share",
        "Customer acquisition cost reduced by 15% through digital channels",
    ]
    for item in items2:
        p = body2.add_paragraph()
        p.text = item
        p.level = 1

    # ---- Slide 3: Key Metrics ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf3 = title3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Key Metrics"
    p3.runs[0].font.size = Pt(32)
    p3.runs[0].font.bold = True
    p3.runs[0].font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # Add a simple metrics table
    tbl_shape = slide3.shapes.add_table(5, 3, Inches(1), Inches(1.5), Inches(10), Inches(3.5))
    tbl = tbl_shape.table
    headers = ["Metric", "Q4 2025", "Target Q4 2026"]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    data_rows = [
        ["Annual Recurring Revenue", "$18.7M", "$27.5M"],
        ["Net Revenue Retention", "112%", "120%"],
        ["Monthly Active Users", "245,000", "400,000"],
        ["Customer Satisfaction (NPS)", "72", "80"],
    ]
    for r, row_data in enumerate(data_rows, 1):
        for c, val in enumerate(row_data):
            tbl.cell(r, c).text = val

    # ---- Slide 4: Our Direction (with arrow shape) ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf4 = title4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Our Direction"
    p4.runs[0].font.size = Pt(32)
    p4.runs[0].font.bold = True
    p4.runs[0].font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # Right-pointing arrow at center of slide (~6in from left, ~3.5in from top)
    arrow = slide4.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(4.5),    # left position: center area
        Inches(2.75),   # top position: roughly center vertically
        Inches(3.0),    # width
        Inches(1.5),    # height
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0x2E, 0x75, 0xB6)
    arrow.line.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

    # ---- Slide 5: Implementation Timeline ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf5 = title5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Implementation Timeline"
    p5.runs[0].font.size = Pt(32)
    p5.runs[0].font.bold = True
    p5.runs[0].font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    phases = [
        ("Phase 1 (Q1 2026)", "Foundation — Hire regional leads, establish partnerships in APAC"),
        ("Phase 2 (Q2 2026)", "Expansion — Launch localized products in 3 new markets"),
        ("Phase 3 (Q3 2026)", "Optimization — Refine pricing, scale marketing spend by 40%"),
        ("Phase 4 (Q4 2026)", "Consolidation — Integrate acquisitions, target 120% NRR"),
    ]
    for i, (phase, desc) in enumerate(phases):
        tb = slide5.shapes.add_textbox(Inches(1), Inches(1.5 + i * 1.3), Inches(10), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p_phase = tf.paragraphs[0]
        p_phase.text = phase
        p_phase.runs[0].font.bold = True
        p_phase.runs[0].font.size = Pt(18)
        p_phase.runs[0].font.color.rgb = RGBColor(0x2E, 0x75, 0xB6)
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.runs[0].font.size = Pt(14)

    # ---- Slide 6: Next Steps ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf6 = title6.text_frame
    p6 = tf6.paragraphs[0]
    p6.text = "Next Steps"
    p6.runs[0].font.size = Pt(32)
    p6.runs[0].font.bold = True
    p6.runs[0].font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    steps = [
        "Finalize regional hiring plan by end of January",
        "Complete partner due diligence for APAC expansion",
        "Present revised pricing model to leadership team",
        "Schedule quarterly business reviews with top 20 accounts",
        "Align product roadmap with expansion requirements",
    ]
    for i, step in enumerate(steps):
        tb = slide6.shapes.add_textbox(Inches(1), Inches(1.5 + i * 0.9), Inches(10), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p_step = tf.paragraphs[0]
        run = p_step.add_run()
        run.text = f"{i+1}. {step}"
        run.font.size = Pt(16)

    # No animations applied — this is the initial state
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
