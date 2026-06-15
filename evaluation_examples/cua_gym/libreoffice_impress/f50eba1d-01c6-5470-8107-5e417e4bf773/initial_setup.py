"""
Initial Setup: Create a presentation with 6 slides using default layout.
Task ID: impress_rp_039
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_039'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard widescreen 13.333 x 7.5 inches
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Dual Column Solutions"
    if slide1.placeholders[1]:
        slide1.placeholders[1].text = "Strategic Planning Workshop 2025"

    # --- Slide 2: Company Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Company Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Founded in 2018, Dual Column Solutions provides enterprise consulting services"
    p2 = body2.add_paragraph()
    p2.text = "Headquartered in San Francisco with offices in London, Tokyo, and Sydney"
    p2.level = 0
    p3 = body2.add_paragraph()
    p3.text = "Over 450 employees across 12 countries"
    p3.level = 0
    p4 = body2.add_paragraph()
    p4.text = "Annual revenue exceeding $85 million in FY2024"
    p4.level = 0

    # --- Slide 3: Market Analysis ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Market Analysis"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "The enterprise consulting market grew 12.3% in 2024"
    items3 = [
        "Digital transformation services lead growth at 18.7% CAGR",
        "AI-driven analytics adoption increased by 34% year-over-year",
        "Cloud migration consulting remains the largest segment ($42B)",
        "Sustainability consulting emerging as fastest-growing niche",
    ]
    for item in items3:
        p = body3.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 4: Service Offerings ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Service Offerings"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Core Services"
    services = [
        "Strategic Planning & Business Architecture",
        "Technology Roadmap Development",
        "Change Management & Organizational Design",
        "Data Analytics & Business Intelligence",
        "Cloud Infrastructure Modernization",
        "Cybersecurity Assessment & Compliance",
    ]
    for svc in services:
        p = body4.add_paragraph()
        p.text = svc
        p.level = 1

    # --- Slide 5: Financial Highlights ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    slide5.shapes.title if hasattr(slide5.shapes, 'title') and slide5.shapes.title else None
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Financial Highlights - FY2024"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # Add a table with financial data
    table_shape = slide5.shapes.add_table(5, 4, Inches(0.5), Inches(1.5), Inches(12), Inches(4))
    table = table_shape.table
    headers = ["Quarter", "Revenue ($M)", "Profit ($M)", "Growth (%)"]
    data_rows = [
        ["Q1 2024", "$19.2", "$3.8", "11.4%"],
        ["Q2 2024", "$21.5", "$4.3", "14.2%"],
        ["Q3 2024", "$22.8", "$4.6", "12.8%"],
        ["Q4 2024", "$23.1", "$4.9", "15.1%"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row_data in enumerate(data_rows, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 6: Next Steps ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Next Steps & Action Items"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Immediate Priorities (Q1 2025)"
    actions = [
        "Complete Phase 2 of digital transformation initiative",
        "Expand APAC operations with new Singapore hub",
        "Launch AI-powered analytics platform for clients",
        "Finalize partnership with CloudScale Technologies",
        "Begin sustainability certification process",
    ]
    for action in actions:
        p = body6.add_paragraph()
        p.text = action
        p.level = 1

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
