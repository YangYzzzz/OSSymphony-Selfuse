"""
Reward Script: Create a slide master with two column text layout and dark sidebar
Task ID: impress_rp_039
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Dark sidebar rectangle (~2in wide, full height, #2C3E50 fill) in a layout
  Component 2 (0.25): Picture/image placeholder inside sidebar area (logo placeholder)
  Component 3 (0.20): Title placeholder on the right side (left >= ~2in)
  Component 4 (0.25): Two body text placeholders side by side on right side
"""

import os
import sys

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_039'


def persist_app_state(domain):
    """Save any unsaved LibreOffice changes before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print(f"PERSIST: ctrl+s sent for {domain}")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Inches, Emu
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Gather all layouts from master 0
    try:
        master = prs.slide_masters[0]
        layouts = list(master.slide_layouts)
        num_layouts = len(layouts)
        print(f"INFO: Found {num_layouts} layouts in master 0")
    except Exception as e:
        print(f"CRITICAL: Cannot access slide master: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Find the custom layout: search all layouts for one that has a non-placeholder
    # shape with solid #2C3E50 fill (the dark sidebar). This is the task-introduced change.
    custom_layout = None
    sidebar_shape = None
    for layout in layouts:
        for shape in layout.shapes:
            if not shape.is_placeholder:
                try:
                    fill = shape.fill
                    if fill.type is not None and fill.type == 1:  # SOLID fill
                        color_rgb = str(fill.fore_color.rgb).upper()
                        if color_rgb == "2C3E50":
                            custom_layout = layout
                            sidebar_shape = shape
                            break
                except Exception:
                    pass
        if custom_layout is not None:
            break

    # Component 1: Dark sidebar rectangle (0.30 points)
    # Must be approximately 2 inches wide, full slide height, filled #2C3E50, on the left
    try:
        if sidebar_shape is not None:
            sidebar_w_in = sidebar_shape.width / 914400.0
            sidebar_h_in = sidebar_shape.height / 914400.0
            slide_h_in = prs.slide_height / 914400.0
            sidebar_left_in = sidebar_shape.left / 914400.0

            checks_passed = 0
            # Width approximately 2 inches (tolerance 0.3in)
            if abs(sidebar_w_in - 2.0) <= 0.3:
                checks_passed += 1
            else:
                print(f"FAIL: Component 1 sub — sidebar width {sidebar_w_in:.2f}in, expected ~2.0in")

            # Full height (at least 90% of slide height)
            if sidebar_h_in >= slide_h_in * 0.9:
                checks_passed += 1
            else:
                print(f"FAIL: Component 1 sub — sidebar height {sidebar_h_in:.2f}in, slide height {slide_h_in:.2f}in")

            # On the left side (left position near 0)
            if sidebar_left_in <= 0.5:
                checks_passed += 1
            else:
                print(f"FAIL: Component 1 sub — sidebar left {sidebar_left_in:.2f}in, expected near 0")

            if checks_passed == 3:
                print(f"PASS: Component 1 — Sidebar: {sidebar_w_in:.2f}in wide, {sidebar_h_in:.2f}in tall, #2C3E50, left={sidebar_left_in:.2f}in (0.30 pts)")
                total_score += 0.30
            else:
                print(f"FAIL: Component 1 — Sidebar partially correct ({checks_passed}/3 sub-checks)")
        else:
            print("FAIL: Component 1 — No shape with #2C3E50 fill found in any layout")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Components 2-4 require the custom layout
    if custom_layout is None:
        print("FAIL: Components 2-4 — No custom layout with sidebar found, skipping")
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Gather placeholders from the custom layout
    placeholders = list(custom_layout.placeholders)
    print(f"INFO: Custom layout '{custom_layout.name}' has {len(placeholders)} placeholders")
    for ph in placeholders:
        pf = ph.placeholder_format
        print(f"  ph idx={pf.idx}, type={pf.type}, name={ph.name}, left={ph.left/914400:.2f}in, top={ph.top/914400:.2f}in, w={ph.width/914400:.2f}in, h={ph.height/914400:.2f}in")

    # Component 2: Picture/image placeholder inside sidebar area (0.25 points)
    # The logo placeholder should be a PICTURE type placeholder positioned within the sidebar
    try:
        picture_phs = []
        for ph in placeholders:
            # PICTURE type is 18
            if ph.placeholder_format.type == 18:
                picture_phs.append(ph)

        found_logo_in_sidebar = False
        for ph in picture_phs:
            ph_left_in = ph.left / 914400.0
            ph_right_in = (ph.left + ph.width) / 914400.0
            # Check if it's within the sidebar area (right edge < 2.5 inches)
            if ph_right_in <= 2.5:
                found_logo_in_sidebar = True
                print(f"PASS: Component 2 — Logo/picture placeholder at left={ph_left_in:.2f}in, right={ph_right_in:.2f}in inside sidebar (0.25 pts)")
                total_score += 0.25
                break

        if not found_logo_in_sidebar:
            if picture_phs:
                print("FAIL: Component 2 — Picture placeholder(s) found but not inside sidebar area")
            else:
                print("FAIL: Component 2 — No picture placeholder found in custom layout")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Title placeholder on the right side (0.20 points)
    # Must be a TITLE type placeholder with left position >= ~2 inches (right of sidebar)
    try:
        title_phs = []
        for ph in placeholders:
            if ph.placeholder_format.type == 1:  # TITLE
                title_phs.append(ph)

        found_title_right = False
        for ph in title_phs:
            ph_left_in = ph.left / 914400.0
            if ph_left_in >= 1.8:  # tolerance: at least 1.8in from left
                found_title_right = True
                print(f"PASS: Component 3 — Title placeholder at left={ph_left_in:.2f}in, width={ph.width/914400:.2f}in (0.20 pts)")
                total_score += 0.20
                break

        if not found_title_right:
            if title_phs:
                for ph in title_phs:
                    print(f"FAIL: Component 3 — Title placeholder at left={ph.left/914400:.2f}in, expected >= 1.8in")
            else:
                print("FAIL: Component 3 — No title placeholder found in custom layout")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Two body text placeholders side by side on right side (0.25 points)
    # Must have two BODY type placeholders, both with left >= ~2in, arranged horizontally
    try:
        body_phs = []
        for ph in placeholders:
            if ph.placeholder_format.type == 2:  # BODY
                ph_left_in = ph.left / 914400.0
                if ph_left_in >= 1.8:  # on the right side
                    body_phs.append(ph)

        if len(body_phs) >= 2:
            # Check they are side by side (similar top, different left positions)
            body_phs_sorted = sorted(body_phs, key=lambda p: p.left)
            ph1 = body_phs_sorted[0]
            ph2 = body_phs_sorted[1]

            top_diff_in = abs(ph1.top - ph2.top) / 914400.0
            left_diff_in = abs(ph1.left - ph2.left) / 914400.0

            # Side by side: similar top (within 0.5in) and different left (at least 1.5in apart)
            if top_diff_in <= 0.5 and left_diff_in >= 1.5:
                # Check approximately equal widths (within 30% of each other)
                w1 = ph1.width / 914400.0
                w2 = ph2.width / 914400.0
                width_ratio = min(w1, w2) / max(w1, w2) if max(w1, w2) > 0 else 0
                if width_ratio >= 0.7:
                    print(f"PASS: Component 4 — Two body placeholders side by side: left1={ph1.left/914400:.2f}in w={w1:.2f}in, left2={ph2.left/914400:.2f}in w={w2:.2f}in (0.25 pts)")
                    total_score += 0.25
                else:
                    print(f"FAIL: Component 4 — Body placeholders have unequal widths: {w1:.2f}in vs {w2:.2f}in (ratio {width_ratio:.2f})")
            else:
                print(f"FAIL: Component 4 — Body placeholders not side by side: top_diff={top_diff_in:.2f}in, left_diff={left_diff_in:.2f}in")
        else:
            print(f"FAIL: Component 4 — Found {len(body_phs)} body placeholder(s) on right side, expected 2")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
