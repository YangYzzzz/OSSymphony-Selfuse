"""
Reward Script: Reconstruct broken SmartArt diagrams on slides 5 and 6
Task ID: impress_fix_023
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 5 has hierarchy connectors (CEO->VPs)         — 0.25 pts
  Component 2: Slide 5 has hierarchy connectors (VPs->Directors)   — 0.35 pts
  Component 3: Slide 6 has arrow connectors between 5 steps        — 0.25 pts
  Component 4: Slide 6 arrows have arrowheads                      — 0.15 pts
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_023'

# XML namespaces for OOXML
NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}


def count_connectors_on_slide(pptx_path, slide_num):
    """Count cxnSp (connector shape) elements on a given slide (1-based)."""
    count = 0
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            slide_xml = f'ppt/slides/slide{slide_num}.xml'
            if slide_xml not in zf.namelist():
                return 0
            with zf.open(slide_xml) as f:
                root = ET.parse(f).getroot()
                cxn_shapes = root.findall('.//p:cSld/p:spTree/p:cxnSp', NS)
                count = len(cxn_shapes)
    except Exception as e:
        print(f"ERROR: Could not parse slide {slide_num} XML: {e}")
    return count


def count_line_shapes_on_slide(pptx_path, slide_num):
    """Count LINE (type 9) shapes via python-pptx as a fallback."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(pptx_path)
        if slide_num - 1 >= len(prs.slides):
            return 0
        slide = prs.slides[slide_num - 1]
        line_count = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                line_count += 1
        return line_count
    except Exception as e:
        print(f"ERROR: Could not count line shapes on slide {slide_num}: {e}")
        return 0


def count_arrows_on_slide(pptx_path, slide_num):
    """Count connector shapes on a slide that have a triangle tailEnd (arrowhead)."""
    count = 0
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            slide_xml = f'ppt/slides/slide{slide_num}.xml'
            if slide_xml not in zf.namelist():
                return 0
            with zf.open(slide_xml) as f:
                root = ET.parse(f).getroot()
                # Check cxnSp connectors for arrowheads
                for sp in root.findall('.//p:cSld/p:spTree/p:cxnSp', NS):
                    ln = sp.find('.//a:ln', NS)
                    if ln is not None:
                        tail = ln.find('a:tailEnd', NS)
                        if tail is not None and tail.get('type') in ('triangle', 'arrow', 'stealth', 'open'):
                            count += 1
                # Also check sp shapes that are lines with arrowheads
                for sp in root.findall('.//p:cSld/p:spTree/p:sp', NS):
                    nv = sp.find('p:nvSpPr/p:cNvPr', NS)
                    name = nv.get('name', '') if nv is not None else ''
                    if 'arrow' in name.lower() or 'connector' in name.lower():
                        ln = sp.find('.//a:ln', NS)
                        if ln is not None:
                            tail = ln.find('a:tailEnd', NS)
                            head = ln.find('a:headEnd', NS)
                            tail_is_arrow = (tail is not None and tail.get('type') in ('triangle', 'arrow', 'stealth', 'open'))
                            head_is_arrow = (head is not None and head.get('type') in ('triangle', 'arrow', 'stealth', 'open'))
                            if tail_is_arrow or head_is_arrow:
                                count += 1
    except Exception as e:
        print(f"ERROR: Could not check arrows on slide {slide_num}: {e}")
    return count


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    if not os.path.exists(file_path):
        print(f"CRITICAL: File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file is a valid pptx with at least 6 slides
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        num_slides = len(prs.slides)
        if num_slides < 6:
            print(f"CRITICAL: Expected at least 6 slides, found {num_slides}")
            print("REWARD: 0.0")
            return 0.0
        print(f"PRECONDITION: File has {num_slides} slides (>= 6 required)")
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide 5 has hierarchy connectors from CEO to VPs (0.25 points)
    # In the golden state, slide 5 should have at least 3 connectors (CEO -> 3 VPs)
    # Total expected: 9 connectors (3 for CEO->VPs, 6 for VPs->Directors)
    # This component checks for at least 3 connectors (CEO->VP level)
    try:
        slide5_cxn = count_connectors_on_slide(file_path, 5)
        slide5_lines = count_line_shapes_on_slide(file_path, 5)
        # Use whichever count is higher (connectors can be cxnSp or sp LINE)
        slide5_connector_count = max(slide5_cxn, slide5_lines)
        print(f"INFO: Slide 5 connectors: {slide5_cxn} cxnSp, {slide5_lines} LINE shapes, effective={slide5_connector_count}")

        if slide5_connector_count >= 3:
            print(f"PASS: Component 1 -- Slide 5 has >= 3 connectors (CEO->VPs hierarchy) ({slide5_connector_count} found) (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 -- Expected >= 3 connectors on slide 5 for CEO->VP links, found {slide5_connector_count}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Slide 5 has full hierarchy connectors (VPs->Directors too) (0.35 points)
    # Expected: 9 total connectors (3 CEO->VP + 6 VP->Dir)
    # Award partial: 6+ connectors = 0.2, 9+ connectors = 0.35
    try:
        if slide5_connector_count >= 9:
            print(f"PASS: Component 2 -- Slide 5 has >= 9 connectors (full org chart hierarchy) (0.35 pts)")
            total_score += 0.35
        elif slide5_connector_count >= 6:
            print(f"PARTIAL: Component 2 -- Slide 5 has {slide5_connector_count} connectors (expected >= 9 for full hierarchy) (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 2 -- Slide 5 has only {slide5_connector_count} connectors, need >= 6 for VP->Dir links")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Slide 6 has arrow connectors between 5 steps (0.25 points)
    # Expected: 4 connectors between 5 process steps
    try:
        slide6_cxn = count_connectors_on_slide(file_path, 6)
        slide6_lines = count_line_shapes_on_slide(file_path, 6)
        slide6_connector_count = max(slide6_cxn, slide6_lines)
        print(f"INFO: Slide 6 connectors: {slide6_cxn} cxnSp, {slide6_lines} LINE shapes, effective={slide6_connector_count}")

        if slide6_connector_count >= 4:
            print(f"PASS: Component 3 -- Slide 6 has >= 4 connectors (process flow connections) (0.25 pts)")
            total_score += 0.25
        elif slide6_connector_count >= 2:
            print(f"PARTIAL: Component 3 -- Slide 6 has {slide6_connector_count} connectors (expected >= 4) (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 3 -- Expected >= 4 connectors on slide 6, found {slide6_connector_count}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Slide 6 arrows have arrowheads (0.15 points)
    # The process flow should use arrows (with triangle/arrow tail ends), not plain lines
    try:
        arrow_count = count_arrows_on_slide(file_path, 6)
        print(f"INFO: Slide 6 arrows with arrowheads: {arrow_count}")

        if arrow_count >= 4:
            print(f"PASS: Component 4 -- Slide 6 has >= 4 arrows with arrowheads (0.15 pts)")
            total_score += 0.15
        elif arrow_count >= 2:
            print(f"PARTIAL: Component 4 -- Slide 6 has {arrow_count} arrows with arrowheads (expected >= 4) (0.07 pts)")
            total_score += 0.07
        else:
            print(f"FAIL: Component 4 -- Expected >= 4 arrows with arrowheads on slide 6, found {arrow_count}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
