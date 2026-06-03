"""
Initial Setup: Broken SmartArt org chart and process flow after PowerPoint import
Task ID: impress_fix_023
Domain: libreoffice_impress

Creates a 6-slide corporate presentation where:
- Slide 5 has scattered rectangles (broken org chart) with NO connectors
- Slide 6 has 5 rounded rectangles (broken process flow) with NO arrows
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_023'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_shape(shape, text, font_size=Pt(12), bold=False, color=None, alignment=PP_ALIGN.CENTER):
    """Helper to add formatted text to a shape."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color scheme
    dark_blue = RGBColor(0x1B, 0x3A, 0x5C)
    medium_blue = RGBColor(0x2E, 0x75, 0xB6)
    light_blue = RGBColor(0x9D, 0xC3, 0xE6)
    white = RGBColor(0xFF, 0xFF, 0xFF)
    dark_gray = RGBColor(0x33, 0x33, 0x33)
    accent_green = RGBColor(0x54, 0x8B, 0x54)

    # ─── Slide 1: Title Slide ───
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = dark_blue

    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    add_text_to_shape(title_box, "Meridian Technologies Inc.", Pt(40), bold=True, color=white)
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
    add_text_to_shape(subtitle_box, "Organizational Structure & Process Overview", Pt(22), color=light_blue)
    date_box = slide1.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(0.6))
    add_text_to_shape(date_box, "Q1 2026 — Confidential", Pt(14), color=RGBColor(0x88, 0x88, 0x88))

    # ─── Slide 2: Agenda ───
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    header2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    add_text_to_shape(header2, "Agenda", Pt(32), bold=True, color=dark_blue, alignment=PP_ALIGN.LEFT)

    agenda_items = [
        "1.  Company Mission & Vision",
        "2.  Financial Highlights — FY2025",
        "3.  Regional Performance Breakdown",
        "4.  Organizational Chart",
        "5.  New Product Development Process",
        "6.  Strategic Priorities for 2026",
    ]
    agenda_box = slide2.shapes.add_textbox(Inches(1), Inches(1.5), Inches(10), Inches(5))
    tf = agenda_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(agenda_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.size = Pt(20)
        run.font.color.rgb = dark_gray
        p.space_after = Pt(14)

    # ─── Slide 3: Company Mission ───
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    header3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    add_text_to_shape(header3, "Our Mission & Vision", Pt(32), bold=True, color=dark_blue, alignment=PP_ALIGN.LEFT)

    mission_box = slide3.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(2))
    tf = mission_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "Meridian Technologies empowers enterprises worldwide through innovative "
        "cloud-native solutions. Since 2008, we have grown from a 12-person startup "
        "in Austin, TX to a global workforce of 4,200+ employees serving Fortune 500 clients "
        "across 28 countries."
    )
    run.font.size = Pt(16)
    run.font.color.rgb = dark_gray

    vision_box = slide3.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(2))
    tf2 = vision_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = (
        "Vision 2030: Become the leading provider of AI-augmented enterprise automation, "
        "achieving $2B ARR while maintaining a 95% customer retention rate."
    )
    run2.font.size = Pt(16)
    run2.font.italic = True
    run2.font.color.rgb = medium_blue

    # ─── Slide 4: Financial Highlights ───
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    header4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    add_text_to_shape(header4, "Financial Highlights — FY2025", Pt(32), bold=True, color=dark_blue, alignment=PP_ALIGN.LEFT)

    # Simple KPI boxes
    kpi_data = [
        ("Revenue", "$847M", "+18% YoY"),
        ("EBITDA", "$203M", "24% margin"),
        ("Headcount", "4,217", "+312 hires"),
        ("NPS", "72", "+5 pts"),
    ]
    for i, (label, value, note) in enumerate(kpi_data):
        x = Inches(0.8 + i * 3.1)
        box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(2.8), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
        box.line.color.rgb = medium_blue
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = value
        run2.font.size = Pt(28)
        run2.font.bold = True
        run2.font.color.rgb = dark_blue

        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        run3 = p3.add_run()
        run3.text = note
        run3.font.size = Pt(12)
        run3.font.color.rgb = accent_green

    # ─── Slide 5: BROKEN Org Chart (scattered rectangles, NO connectors) ───
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    header5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    add_text_to_shape(header5, "Organizational Structure", Pt(32), bold=True, color=dark_blue, alignment=PP_ALIGN.LEFT)

    # CEO box - top center
    ceo_box = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(1.3), Inches(2.8), Inches(0.9))
    ceo_box.fill.solid()
    ceo_box.fill.fore_color.rgb = dark_blue
    ceo_box.line.color.rgb = dark_blue
    add_text_to_shape(ceo_box, "Elena Rodriguez\nCEO", Pt(13), bold=True, color=white)

    # VP boxes - middle row (scattered slightly off-grid to look broken)
    vp_data = [
        ("James Chen\nVP Engineering", Inches(1.0), Inches(3.0)),
        ("Sarah Mitchell\nVP Sales", Inches(5.0), Inches(2.9)),
        ("David Park\nVP Operations", Inches(9.2), Inches(3.1)),
    ]
    vp_boxes = []
    for name, x, y in vp_data:
        vp = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.8), Inches(0.85))
        vp.fill.solid()
        vp.fill.fore_color.rgb = medium_blue
        vp.line.color.rgb = medium_blue
        add_text_to_shape(vp, name, Pt(11), bold=True, color=white)
        vp_boxes.append(vp)

    # Director boxes - bottom row (scattered)
    dir_data = [
        ("Lisa Wang\nDir. Frontend", Inches(0.3), Inches(5.2)),
        ("Tom Harris\nDir. Backend", Inches(2.3), Inches(5.3)),
        ("Amy Foster\nDir. Enterprise", Inches(4.0), Inches(5.1)),
        ("Rick Nguyen\nDir. SMB", Inches(6.2), Inches(5.3)),
        ("Karen Lee\nDir. Logistics", Inches(8.5), Inches(5.1)),
        ("Marcus Brown\nDir. Quality", Inches(10.5), Inches(5.2)),
    ]
    for name, x, y in dir_data:
        d = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.3), Inches(0.8))
        d.fill.solid()
        d.fill.fore_color.rgb = light_blue
        d.line.color.rgb = medium_blue
        add_text_to_shape(d, name, Pt(10), color=dark_blue)

    # NOTE: NO connector lines — this is the "broken" state

    # ─── Slide 6: BROKEN Process Flow (rounded rects, NO arrows) ───
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    header6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    add_text_to_shape(header6, "New Product Development Process", Pt(32), bold=True, color=dark_blue, alignment=PP_ALIGN.LEFT)

    steps = [
        ("Step 1\nIdeation &\nResearch", RGBColor(0x2E, 0x75, 0xB6)),
        ("Step 2\nPrototype\nDevelopment", RGBColor(0x40, 0x8E, 0xC6)),
        ("Step 3\nUser Testing\n& Feedback", RGBColor(0x54, 0x8B, 0x54)),
        ("Step 4\nRefinement\n& QA", RGBColor(0xD4, 0x8B, 0x2C)),
        ("Step 5\nLaunch &\nMonitoring", RGBColor(0xC0, 0x39, 0x2B)),
    ]

    for i, (label, color) in enumerate(steps):
        x = Inches(0.5 + i * 2.55)
        y = Inches(2.8)
        box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.2), Inches(2.0))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = color
        add_text_to_shape(box, label, Pt(14), bold=True, color=white)

    # NOTE: NO arrows between steps — this is the "broken" state

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
