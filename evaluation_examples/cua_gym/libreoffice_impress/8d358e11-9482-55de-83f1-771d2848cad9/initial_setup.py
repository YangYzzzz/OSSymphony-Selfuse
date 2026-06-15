"""
Initial Setup: Create risk_assessment.pptx with 5 slides, slide 4 titled 'Risk Heat Map' but empty.
Task ID: impress_gf5_047
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_047'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 alignment=PP_ALIGN.LEFT, color=None):
    """Helper to add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullet_slide(prs, title_text, bullets):
    """Add a slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(16)
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Project Risk Assessment"
    slide1.placeholders[1].text = "Q2 2025 Portfolio Review\nPrepared by: Enterprise Risk Management Team"

    # --- Slide 2: Risk Overview ---
    add_bullet_slide(prs, "Risk Assessment Methodology", [
        "Qualitative risk analysis using Likelihood x Impact matrix",
        "Risk ratings determined by cross-functional review board",
        "Assessment covers financial, operational, and strategic risks",
        "Updated quarterly with monthly exception reporting",
        "Aligned with ISO 31000 risk management framework",
        "Stakeholder interviews conducted across 12 business units",
    ])

    # --- Slide 3: Key Risk Areas ---
    add_bullet_slide(prs, "Key Risk Areas Identified", [
        "Supply chain disruption - semiconductor shortages persist",
        "Cybersecurity threats - ransomware incidents up 34% YoY",
        "Regulatory compliance - new data privacy regulations in EU and APAC",
        "Talent retention - attrition rate at 18.5% in engineering",
        "Foreign exchange exposure - USD/EUR volatility impacts margin",
        "Climate-related risks - physical asset exposure in coastal regions",
        "Third-party vendor concentration - top 3 vendors = 62% of spend",
    ])

    # --- Slide 4: Risk Heat Map (EMPTY - agent task) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Just add the title text box at the top
    add_text_box(slide4, Inches(0.5), Inches(0.2), Inches(9), Inches(0.7),
                 "Risk Heat Map", font_size=28, bold=True,
                 alignment=PP_ALIGN.CENTER,
                 color=RGBColor(0x1F, 0x38, 0x64))

    # --- Slide 5: Mitigation Strategies ---
    add_bullet_slide(prs, "Mitigation Strategies", [
        "Diversify supplier base - onboard 3 alternative semiconductor vendors by Q3",
        "Implement zero-trust security architecture across all business units",
        "Establish regulatory monitoring dashboard for 15 key jurisdictions",
        "Launch retention program with competitive equity packages for top performers",
        "Deploy natural hedging strategy and increase FX forward coverage to 80%",
        "Complete climate risk scenario analysis for all facilities by end of year",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
