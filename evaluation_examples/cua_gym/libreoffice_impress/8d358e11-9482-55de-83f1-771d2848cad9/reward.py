"""
Reward Script: Risk Heat Map on Slide 4 of risk_assessment.pptx
Task ID: impress_gf5_047
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30) - 25 colored rectangles exist on slide 4 (5x5 grid)
  Component 2 (0.35) - Rectangle risk scores match Likelihood x Impact values
  Component 3 (0.20) - Colors match risk thresholds (1-4=green, 5-9=yellow, 10-14=orange, 15-25=red)
  Component 4 (0.15) - Axis labels (Likelihood, Impact) and title present
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_047'


def persist_app_state(domain: str):
    """Try to save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(1.0)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)

    # Collect all AUTO_SHAPE rectangles that contain numeric text (grid cells)
    grid_rects = []
    legend_rects = []
    all_text_content = []

    for shape in slide.shapes:
        # Collect all text on the slide for label checks
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    all_text_content.append(text.lower())

        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            text = ''
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()

            # Get fill color
            fill_color = None
            try:
                if shape.fill.type is not None:
                    fill_color = str(shape.fill.fore_color.rgb).upper()
            except Exception:
                pass

            if text and text.isdigit():
                grid_rects.append({
                    'text': int(text),
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                    'color': fill_color,
                })
            elif not text and fill_color:
                # Legend color swatch (no text, has color)
                legend_rects.append(fill_color)

    print(f"INFO: Found {len(grid_rects)} grid rectangles with numeric text")
    print(f"INFO: Found {len(legend_rects)} legend color swatches")

    # =========================================================================
    # Component 1: 25 colored rectangles on slide 4 forming a 5x5 grid (0.30)
    # =========================================================================
    try:
        if len(grid_rects) >= 25:
            print(f"PASS: Component 1 — Found {len(grid_rects)} grid rectangles (>= 25) (0.30 pts)")
            total_score += 0.30
        elif len(grid_rects) >= 15:
            partial = 0.30 * (len(grid_rects) / 25.0)
            print(f"PARTIAL: Component 1 — Found {len(grid_rects)}/25 grid rectangles ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — Only {len(grid_rects)} grid rectangles found, need 25")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # =========================================================================
    # Component 2: Risk scores match Likelihood x Impact (0.35)
    # The grid should contain all 25 products: L*I for L in 1..5 and I in 1..5
    # Expected values: {1,2,3,4,5,6,8,9,10,12,15,16,20,25} with correct multiplicities
    # =========================================================================
    try:
        expected_scores = []
        for likelihood in range(1, 6):
            for impact in range(1, 6):
                expected_scores.append(likelihood * impact)
        expected_scores_sorted = sorted(expected_scores)

        actual_scores = sorted([r['text'] for r in grid_rects])

        if len(grid_rects) >= 25:
            # Compare the 25 values
            actual_25 = actual_scores[:25] if len(actual_scores) >= 25 else actual_scores
            matching = sum(1 for a, e in zip(actual_25, expected_scores_sorted) if a == e)
            ratio = matching / 25.0
            pts = 0.35 * ratio
            if ratio >= 0.95:
                print(f"PASS: Component 2 — {matching}/25 risk scores match L*I values (0.35 pts)")
                total_score += 0.35
            elif ratio > 0.0:
                print(f"PARTIAL: Component 2 — {matching}/25 risk scores match ({pts:.2f} pts)")
                total_score += pts
            else:
                print(f"FAIL: Component 2 — No risk scores match expected L*I values")
                print(f"  Expected: {expected_scores_sorted}")
                print(f"  Actual:   {actual_25}")
        else:
            print(f"FAIL: Component 2 — Not enough rectangles ({len(grid_rects)}) to verify scores")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # =========================================================================
    # Component 3: Colors match risk thresholds (0.20)
    # 1-4 = green, 5-9 = yellow, 10-14 = orange, 15-25 = red
    # We check each rectangle's color against its score
    # =========================================================================
    try:
        def get_expected_color_category(score):
            """Return the expected color category for a risk score."""
            if score <= 4:
                return 'green'
            elif score <= 9:
                return 'yellow'
            elif score <= 14:
                return 'orange'
            else:
                return 'red'

        # Known color values from the golden file, but also accept common variants
        GREEN_COLORS = {'2ECC40', '00FF00', '00B050', '008000', '00A550', '22B14C',
                        '92D050', '00B04F', '28A745', '00FF7F', '32CD32', '2ECC40'}
        YELLOW_COLORS = {'FFDC00', 'FFFF00', 'FFD700', 'FFC000', 'FFBF00', 'F0E68C',
                         'FFD600', 'FFC107', 'FFEB3B'}
        ORANGE_COLORS = {'FF851B', 'FF8C00', 'FFA500', 'FF6600', 'ED7D31', 'FF7F50',
                         'E67E22', 'FF9800', 'F08C00'}
        RED_COLORS = {'FF4136', 'FF0000', 'CC0000', 'FF0033', 'C00000', 'FF3333',
                      'DC3545', 'F44336', 'E74C3C', 'FF1744'}

        def classify_color(hex_color):
            """Classify a hex color string into a risk category."""
            if hex_color is None:
                return None
            hex_color = hex_color.upper().replace('#', '')
            if hex_color in GREEN_COLORS:
                return 'green'
            if hex_color in YELLOW_COLORS:
                return 'yellow'
            if hex_color in ORANGE_COLORS:
                return 'orange'
            if hex_color in RED_COLORS:
                return 'red'
            # Fallback: classify by RGB channel dominance
            try:
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)
                if g > r and g > 100:
                    return 'green'
                elif r > 200 and g > 180 and b < 100:
                    return 'yellow'
                elif r > 200 and g > 80 and g < 180 and b < 80:
                    return 'orange'
                elif r > 200 and g < 100:
                    return 'red'
            except:
                pass
            return None

        if len(grid_rects) >= 25:
            correct_colors = 0
            for rect in grid_rects:
                expected_cat = get_expected_color_category(rect['text'])
                actual_cat = classify_color(rect['color'])
                if expected_cat == actual_cat:
                    correct_colors += 1
                else:
                    print(f"  Color mismatch: score={rect['text']}, expected={expected_cat}, "
                          f"got={actual_cat} (hex={rect['color']})")

            ratio = correct_colors / len(grid_rects)
            pts = 0.20 * ratio
            if ratio >= 0.90:
                print(f"PASS: Component 3 — {correct_colors}/{len(grid_rects)} cells have correct colors (0.20 pts)")
                total_score += 0.20
            elif ratio > 0.0:
                print(f"PARTIAL: Component 3 — {correct_colors}/{len(grid_rects)} correct colors ({pts:.2f} pts)")
                total_score += pts
            else:
                print(f"FAIL: Component 3 — No cells have correct risk threshold colors")
        else:
            print(f"FAIL: Component 3 — Not enough grid rectangles to check colors")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # =========================================================================
    # Component 4: Axis labels present AND grid exists (0.15)
    # These labels only matter if the grid was actually created (task-introduced).
    # The title "Risk Heat Map" is a precondition (exists in initial), so we only
    # score the NEW labels: "Likelihood" and "Impact", gated on grid existence.
    # =========================================================================
    try:
        # Gate: only award points if grid rectangles were actually added
        if len(grid_rects) >= 15:
            all_text_lower = ' '.join(all_text_content)
            sub_score = 0.0

            has_likelihood = 'likelihood' in all_text_lower
            has_impact = 'impact' in all_text_lower

            if has_likelihood:
                sub_score += 0.075
                print(f"  PASS: 'Likelihood' axis label found")
            else:
                print(f"  FAIL: 'Likelihood' axis label not found")

            if has_impact:
                sub_score += 0.075
                print(f"  PASS: 'Impact' axis label found")
            else:
                print(f"  FAIL: 'Impact' axis label not found")

            if sub_score >= 0.15:
                print(f"PASS: Component 4 — Axis labels present with grid (0.15 pts)")
                total_score += sub_score
            elif sub_score > 0.0:
                print(f"PARTIAL: Component 4 — Some labels present ({sub_score:.2f} pts)")
                total_score += sub_score
            else:
                print(f"FAIL: Component 4 — No axis labels found despite grid being present")
        else:
            print(f"FAIL: Component 4 — Grid not present, axis labels not scored")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
