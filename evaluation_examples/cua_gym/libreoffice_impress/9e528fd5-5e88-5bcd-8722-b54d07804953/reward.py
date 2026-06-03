"""
Reward Script: Restructure slide 9 content as a two-column comparison table
Task ID: impress_cross_acad_027
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 9 has a TABLE shape (0.3 pts)
  Component 2: Table has correct 2-column x 5-row dimensions (0.3 pts)
  Component 3: Header row has correct column headers (0.2 pts)
  Component 4: Data rows contain the 4 comparison points (0.2 pts)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user/Documents'
TASK_ID = 'impress_cross_acad_027'

def normalize_text(t):
    """Normalize text for comparison: strip whitespace, lowercase."""
    return t.strip().lower() if t else ''


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Task: In slide 9, the content text box should be replaced with a 2-column
    comparison table. Left column header: 'Classical Mechanics', right column
    header: 'Quantum Mechanics', 4 data rows from the original comparison text.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Verify slide count precondition gate
    if len(prs.slides) < 9:
        print(f"CRITICAL: Expected at least 9 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Get slide 9 (0-indexed: index 8)
    slide9 = prs.slides[8]

    # Component 1: Slide 9 has a TABLE shape (0.3 points)
    # In initial file, slide 9 has no table — only text placeholders.
    # Task requires replacing the text box with a table.
    try:
        table_shapes = [s for s in slide9.shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE]
        if len(table_shapes) >= 1:
            print(f"PASS: Component 1 — Slide 9 has {len(table_shapes)} TABLE shape(s) (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 — Slide 9 has no TABLE shape. Shapes found: {[s.shape_type for s in slide9.shapes]}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Table has correct 2-column x 5-row dimensions (0.3 points)
    # Ground truth: 2-column, 5-row table (1 header row + 4 data rows)
    try:
        table_shapes = [s for s in slide9.shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE]
        if table_shapes:
            table = table_shapes[0].table
            num_rows = len(table.rows)
            num_cols = len(table.columns)
            if num_rows == 5 and num_cols == 2:
                print(f"PASS: Component 2 — Table dimensions are {num_rows} rows x {num_cols} cols (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 2 — Expected 5 rows x 2 cols, found {num_rows} rows x {num_cols} cols")
        else:
            print("FAIL: Component 2 — No table found on slide 9")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Header row has correct column headers (0.2 points)
    # Ground truth: Cell(0,0)="Classical Mechanics", Cell(0,1)="Quantum Mechanics"
    try:
        table_shapes = [s for s in slide9.shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE]
        if table_shapes:
            table = table_shapes[0].table
            if len(table.rows) >= 1 and len(table.columns) >= 2:
                header_left = table.cell(0, 0).text.strip()
                header_right = table.cell(0, 1).text.strip()
                left_ok = normalize_text(header_left) == 'classical mechanics'
                right_ok = normalize_text(header_right) == 'quantum mechanics'
                if left_ok and right_ok:
                    print(f"PASS: Component 3 — Headers correct: '{header_left}' | '{header_right}' (0.2 pts)")
                    total_score += 0.2
                else:
                    print(f"FAIL: Component 3 — Expected headers 'Classical Mechanics'|'Quantum Mechanics', "
                          f"found '{header_left}'|'{header_right}'")
            else:
                print("FAIL: Component 3 — Table too small to check headers")
        else:
            print("FAIL: Component 3 — No table found on slide 9")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Data rows contain the 4 comparison points (0.2 points)
    # Ground truth: 4 data rows with content from original text:
    #   Row 1: Classical "Deterministic" | Quantum "Probabilistic"
    #   Row 2: Classical "Continuous energy" | Quantum "Discrete energy levels"
    #   Row 3: Classical "Particle trajectories" | Quantum "Wave functions"
    #   Row 4: Classical "Macroscale" | Quantum "Nanoscale"
    # We check that the key content terms are present in the table data rows.
    try:
        table_shapes = [s for s in slide9.shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE]
        if table_shapes:
            table = table_shapes[0].table
            if len(table.rows) >= 5 and len(table.columns) >= 2:
                # Extract all data row content (rows 1-4, 0-indexed)
                data_rows = []
                for r in range(1, 5):
                    left = normalize_text(table.cell(r, 0).text)
                    right = normalize_text(table.cell(r, 1).text)
                    data_rows.append((left, right))

                # Check that the key content keywords are present
                # Expected pairs (normalized): left side and right side keywords
                expected_keywords = [
                    ('deterministic', 'probabilistic'),
                    ('continuous', 'discrete'),
                    ('trajectories', 'wave'),
                    ('macroscale', 'nanoscale'),
                ]

                matches = 0
                for exp_left, exp_right in expected_keywords:
                    pair_match = any(exp_left in left and exp_right in right for left, right in data_rows)
                    if pair_match:
                        matches += 1
                    else:
                        print(f"  Missing pair: '{exp_left}' | '{exp_right}'")

                if matches == 4:
                    print(f"PASS: Component 4 — All 4 data rows have correct comparison content (0.2 pts)")
                    total_score += 0.2
                else:
                    print(f"FAIL: Component 4 — Only {matches}/4 data row pairs matched. "
                          f"Data rows: {data_rows}")
            else:
                print("FAIL: Component 4 — Table too small to check 4 data rows")
        else:
            print("FAIL: Component 4 — No table found on slide 9")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against golden file (path on VM)
file_path = f'{WORKDIR}/{TASK_ID}_initial.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
