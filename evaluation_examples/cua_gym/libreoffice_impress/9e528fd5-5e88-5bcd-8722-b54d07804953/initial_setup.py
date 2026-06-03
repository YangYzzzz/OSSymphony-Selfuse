"""
Initial Setup: Physics lecture presentation with 22 slides.
Slide 9 has a text box with comparison content (no table yet).
Task ID: impress_cross_acad_027
Domain: libreoffice_impress
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_cross_acad_027'
OUTPUT = f'{WORKDIR}/Documents/{TASK_ID}_initial.pptx'

# Slide content for a 22-slide physics lecture
SLIDE_CONTENT = [
    # (title, body)
    ("Introduction to Classical and Quantum Physics",
     "A comparative study of two foundational frameworks in physics"),
    ("Course Overview",
     "Part 1: Classical Mechanics\nPart 2: Quantum Mechanics\nPart 3: Bridging the Gap\nPart 4: Modern Applications"),
    ("Historical Context",
     "Newton's Principia Mathematica (1687) established classical mechanics\nPlanck's quantum hypothesis (1900) revolutionized physics\nEinstein's photoelectric effect (1905) confirmed quantum theory"),
    ("Classical Mechanics — Core Principles",
     "Newton's Three Laws of Motion\nConservation of Energy\nConservation of Momentum\nGravitational and Electromagnetic Forces"),
    ("Classical Mechanics — Key Equations",
     "F = ma (Newton's Second Law)\nE_k = ½mv² (Kinetic Energy)\nF = Gm₁m₂/r² (Gravitational Force)\np = mv (Momentum)"),
    ("The Macroscopic World",
     "Classical mechanics excels at describing:\n• Planet orbits and celestial mechanics\n• Structural engineering and bridges\n• Fluid dynamics and aerodynamics\n• Mechanical systems and machines"),
    ("Introduction to Quantum Theory",
     "Energy comes in discrete packets called quanta\nParticles exhibit wave-particle duality\nMeasurement affects quantum states\nProbabilistic rather than deterministic predictions"),
    ("Quantum Mechanics — Key Equations",
     "E = hν (Energy of a photon)\nΔx·Δp ≥ ℏ/2 (Heisenberg Uncertainty)\nHψ = Eψ (Schrödinger Equation)\np = ℏk (De Broglie relation)"),
    ("Comparing Frameworks",
     "Classical: deterministic | Quantum: probabilistic\nClassical: continuous energy | Quantum: discrete energy levels\nClassical: particle trajectories | Quantum: wave functions\nClassical: macroscale | Quantum: nanoscale"),
    ("Wave-Particle Duality",
     "Light behaves as both wave (interference) and particle (photoelectric effect)\nElectrons exhibit diffraction patterns\nDe Broglie wavelength: λ = h/p\nComplementarity principle (Bohr)"),
    ("The Uncertainty Principle",
     "Position and momentum cannot both be precisely known\nΔx·Δp ≥ ℏ/2\nEnergy and time: ΔE·Δt ≥ ℏ/2\nFundamental limit, not a measurement artifact"),
    ("Quantum States and Superposition",
     "Quantum systems exist in superpositions of states\nSchrödinger's cat thought experiment\nMeasurement collapses the wave function\nEntanglement: non-local correlations"),
    ("Atomic Models",
     "Thomson's plum pudding model (1904)\nRutherford's nuclear model (1911)\nBohr's quantized orbits (1913)\nSchrödinger's wave mechanical model (1926)"),
    ("Energy Quantization",
     "Blackbody radiation problem\nPhotoelectric effect (Einstein, 1905)\nAtomic emission spectra\nZero-point energy"),
    ("Transition from Classical to Quantum",
     "Correspondence principle: quantum → classical at large scales\nPlanck's constant ℏ = 1.055 × 10⁻³⁴ J·s\nBohr's limit: n → ∞ gives classical orbits\nDecoherence explains classical appearance"),
    ("Quantum Applications in Technology",
     "Transistors and semiconductors (quantum tunneling)\nLasers (stimulated emission)\nMRI machines (nuclear spin)\nQuantum cryptography and computing"),
    ("Classical Applications in Engineering",
     "Structural analysis (beams, bridges, buildings)\nFluid mechanics (aircraft, ships)\nThermodynamics (engines, turbines)\nOptical instruments (telescopes, microscopes)"),
    ("Limitations of Each Framework",
     "Classical fails at: atomic scales, high speeds (→ relativity), quantum phenomena\nQuantum fails at: macroscopic intuition, computational tractability for large systems\nBoth are approximations of a deeper theory"),
    ("The Path to Unification",
     "Quantum Field Theory unifies QM and special relativity\nStandard Model of particle physics\nGeneral Relativity remains separate\nQuantum gravity remains an open problem"),
    ("Experimental Methods",
     "Classical: macroscopic instruments, direct measurement\nQuantum: particle accelerators, interferometers, spectroscopy\nBoth: statistical analysis of large datasets\nModern: tabletop quantum experiments"),
    ("Summary and Key Takeaways",
     "Classical mechanics governs the macroscale world\nQuantum mechanics governs atoms and subatomic particles\nBoth are essential for modern science and engineering\nUnification remains one of physics' greatest challenges"),
    ("References and Further Reading",
     "Feynman Lectures on Physics (Feynman, Leighton, Sands)\nIntroduction to Quantum Mechanics (Griffiths)\nClassical Mechanics (Goldstein)\nThe Quantum World (French & Taylor)"),
]


def create_initial():
    prs = Presentation()
    # Standard widescreen slide dimensions
    # Default is 10 x 7.5 inches (9144000 x 6858000 EMU)

    for i, (title_text, body_text) in enumerate(SLIDE_CONTENT):
        if i == 0:
            # Title slide
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title_text
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = body_text
        else:
            # Title + Content layout
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title_text
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = body_text

    # Verify slide count
    assert len(prs.slides) == 22, f"Expected 22 slides, got {len(prs.slides)}"

    # Verify slide 9 (index 8) has the comparison text box (not a table)
    slide9 = prs.slides[8]
    assert slide9.shapes.title.text == "Comparing Frameworks", \
        f"Slide 9 title mismatch: {slide9.shapes.title.text}"

    # Make sure output directory exists on VM — handled by script path
    os.makedirs(f'{WORKDIR}/Documents', exist_ok=True)
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'  Slides: {len(prs.slides)}')
    print(f'  Slide 9 title: {prs.slides[8].shapes.title.text}')
    # Print text content of slide 9
    slide9 = prs.slides[8]
    for shape in slide9.shapes:
        if shape.has_text_frame and shape != slide9.shapes.title:
            print(f'  Slide 9 body text: {repr(shape.text_frame.text[:100])}')


create_initial()
