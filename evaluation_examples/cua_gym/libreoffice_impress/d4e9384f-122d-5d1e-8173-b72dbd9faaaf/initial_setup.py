"""
Initial Setup: Create Tech_Architecture.pptx with 8 slides.
Slide 2 has title 'System Architecture' and empty content area.
Task ID: impress_ps_029
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_029'
OUTPUT = f'{WORKDIR}/Tech_Architecture.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_shape(shape, text, font_size=Pt(14), bold=False, color=None, alignment=None):
    """Helper to set text on a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Tech Architecture Overview"
    slide1.placeholders[1].text = "Engineering Team — Q2 2025 Review"

    # --- Slide 2: System Architecture (EMPTY - task target) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title text box manually
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "System Architecture"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x3B, 0x4E)
    # No shapes, connectors, or diagrams — just the title

    # --- Slide 3: Technology Stack ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    tb3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Technology Stack"
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.runs[0]
    r3.font.size = Pt(28)
    r3.font.bold = True
    r3.font.color.rgb = RGBColor(0x2E, 0x3B, 0x4E)

    # Add content bullets
    content3 = slide3.shapes.add_textbox(Inches(1), Inches(1.2), Inches(5), Inches(5))
    tf3c = content3.text_frame
    tf3c.word_wrap = True
    items = [
        ("Frontend", "React 18.2, TypeScript 5.0, Tailwind CSS 3.4"),
        ("Backend", "Node.js 20 LTS, Express 4.18, GraphQL"),
        ("Database", "PostgreSQL 16, Redis 7.2, Elasticsearch 8.11"),
        ("Infrastructure", "AWS EKS, Terraform, ArgoCD"),
        ("Monitoring", "Datadog, PagerDuty, Grafana"),
    ]
    for i, (cat, detail) in enumerate(items):
        if i == 0:
            p = tf3c.paragraphs[0]
        else:
            p = tf3c.add_paragraph()
        p.text = f"{cat}: {detail}"
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.size = Pt(14)

    # --- Slide 4: Deployment Pipeline ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    tb4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf4 = tb4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "CI/CD Pipeline"
    p4.alignment = PP_ALIGN.CENTER
    r4 = p4.runs[0]
    r4.font.size = Pt(28)
    r4.font.bold = True
    r4.font.color.rgb = RGBColor(0x2E, 0x3B, 0x4E)

    stages = [
        "1. Code Push → GitHub Actions triggered",
        "2. Unit Tests + Integration Tests (Jest, Pytest)",
        "3. Static Analysis (ESLint, SonarQube)",
        "4. Container Build (Docker multi-stage)",
        "5. Deploy to Staging (EKS via Helm charts)",
        "6. Smoke Tests + Performance Regression",
        "7. Manual Approval Gate",
        "8. Production Rolling Deploy (canary 10% → 50% → 100%)",
    ]
    content4 = slide4.shapes.add_textbox(Inches(1), Inches(1.2), Inches(10), Inches(5.5))
    tf4c = content4.text_frame
    tf4c.word_wrap = True
    for i, stage in enumerate(stages):
        if i == 0:
            p = tf4c.paragraphs[0]
        else:
            p = tf4c.add_paragraph()
        p.text = stage
        p.space_after = Pt(6)
        run = p.runs[0]
        run.font.size = Pt(14)

    # --- Slide 5: Performance Metrics ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    tb5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf5 = tb5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Performance Metrics — March 2025"
    p5.alignment = PP_ALIGN.CENTER
    r5 = p5.runs[0]
    r5.font.size = Pt(28)
    r5.font.bold = True
    r5.font.color.rgb = RGBColor(0x2E, 0x3B, 0x4E)

    # Table with metrics
    rows, cols = 6, 4
    table_shape = slide5.shapes.add_table(rows, cols, Inches(1.5), Inches(1.5), Inches(10), Inches(4))
    table = table_shape.table
    headers = ["Metric", "Target", "Actual", "Status"]
    data = [
        ["API P99 Latency", "< 200ms", "187ms", "On Track"],
        ["Uptime (SLA)", "99.95%", "99.98%", "Exceeding"],
        ["Error Rate", "< 0.1%", "0.04%", "On Track"],
        ["Avg Response Time", "< 50ms", "38ms", "Exceeding"],
        ["Throughput", "10K req/s", "12.3K req/s", "Exceeding"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(12)
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(11)

    # --- Slide 6: Security Overview ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    tb6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf6 = tb6.text_frame
    p6 = tf6.paragraphs[0]
    p6.text = "Security Posture"
    p6.alignment = PP_ALIGN.CENTER
    r6 = p6.runs[0]
    r6.font.size = Pt(28)
    r6.font.bold = True
    r6.font.color.rgb = RGBColor(0x2E, 0x3B, 0x4E)

    sec_items = [
        "• OAuth 2.0 + OIDC for authentication via Okta",
        "• Role-based access control (RBAC) with 4 permission tiers",
        "• TLS 1.3 enforced on all external endpoints",
        "• Secrets managed via HashiCorp Vault (auto-rotation every 30 days)",
        "• WAF rules updated weekly; last pen test: Feb 2025 — 0 critical findings",
        "• SOC 2 Type II audit completed January 2025",
    ]
    content6 = slide6.shapes.add_textbox(Inches(1), Inches(1.2), Inches(10), Inches(5.5))
    tf6c = content6.text_frame
    tf6c.word_wrap = True
    for i, item in enumerate(sec_items):
        if i == 0:
            p = tf6c.paragraphs[0]
        else:
            p = tf6c.add_paragraph()
        p.text = item
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.size = Pt(14)

    # --- Slide 7: Scaling Strategy ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    tb7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf7 = tb7.text_frame
    p7 = tf7.paragraphs[0]
    p7.text = "Horizontal Scaling Strategy"
    p7.alignment = PP_ALIGN.CENTER
    r7 = p7.runs[0]
    r7.font.size = Pt(28)
    r7.font.bold = True
    r7.font.color.rgb = RGBColor(0x2E, 0x3B, 0x4E)

    scale_items = [
        "Current: 3 EKS node groups (t3.xlarge) across 3 AZs",
        "Auto-scaling: HPA with CPU 70% / Memory 80% thresholds",
        "Database: Read replicas (2x) + connection pooling via PgBouncer",
        "Cache: Redis Cluster mode, 6 shards, 2 replicas each",
        "CDN: CloudFront with 47 edge locations, 98.2% cache hit ratio",
        "Load Balancer: ALB with weighted target groups for canary deploys",
    ]
    content7 = slide7.shapes.add_textbox(Inches(1), Inches(1.2), Inches(10), Inches(5.5))
    tf7c = content7.text_frame
    tf7c.word_wrap = True
    for i, item in enumerate(scale_items):
        if i == 0:
            p = tf7c.paragraphs[0]
        else:
            p = tf7c.add_paragraph()
        p.text = item
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.size = Pt(14)

    # --- Slide 8: Roadmap ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    tb8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf8 = tb8.text_frame
    p8 = tf8.paragraphs[0]
    p8.text = "Q3 2025 Architecture Roadmap"
    p8.alignment = PP_ALIGN.CENTER
    r8 = p8.runs[0]
    r8.font.size = Pt(28)
    r8.font.bold = True
    r8.font.color.rgb = RGBColor(0x2E, 0x3B, 0x4E)

    roadmap = [
        "July: Migrate remaining monolith services to microservices",
        "July: Implement event-driven architecture with Apache Kafka",
        "August: Deploy service mesh (Istio) for inter-service communication",
        "August: Complete GraphQL federation across all domains",
        "September: Launch real-time analytics pipeline (Flink + Druid)",
        "September: Achieve ISO 27001 certification",
    ]
    content8 = slide8.shapes.add_textbox(Inches(1), Inches(1.2), Inches(10), Inches(5.5))
    tf8c = content8.text_frame
    tf8c.word_wrap = True
    for i, item in enumerate(roadmap):
        if i == 0:
            p = tf8c.paragraphs[0]
        else:
            p = tf8c.add_paragraph()
        p.text = item
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.size = Pt(14)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
