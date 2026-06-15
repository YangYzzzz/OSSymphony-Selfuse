"""
Reward Script: System Architecture Overview on Slide 2
Task ID: impress_ps_029
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Three tier rectangles with correct labels on slide 2
  Component 2 (0.40): Six component boxes with correct labels on slide 2
  Component 3 (0.30): Two downward arrow connectors between tiers on slide 2
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_029'


def persist_app_state(domain):
    """Save any unsaved LibreOffice edits via Ctrl+S."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(1.0)
            print("PERSIST: ctrl+s sent for libreoffice_impress")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_all_text_shapes(slide):
    """Recursively get all shapes that have text frames, including inside groups."""
    def extract(shape):
        results = []
        if hasattr(shape, "text") and hasattr(shape, "text_frame"):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.extend(extract(sub))
        return results
    out = []
    for shape in slide.shapes:
        out.extend(extract(shape))
    return out


def _check_arrow_head(ln_element):
    """Check if a line element has an arrow head (triangle/arrow/stealth)."""
    if ln_element is None:
        return False
    for tag_name in ('a:tailEnd', 'a:headEnd'):
        end = ln_element.find(qn(tag_name))
        if end is not None and end.get('type') in ('triangle', 'arrow', 'stealth'):
            return True  # derived from XML attribute check
    return False


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 2")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[1]  # Slide 2 (0-indexed)

    # Collect all shapes on slide 2 by type
    auto_shapes = []  # rectangles / rounded rectangles
    connectors = []   # lines / connectors

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            text = ""
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
            auto_shapes.append({
                'name': shape.name,
                'text': text,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'shape': shape,
            })
        elif shape.shape_type == MSO_SHAPE_TYPE.LINE or shape.shape_type == 9:
            connectors.append({
                'name': shape.name,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'element': shape._element,
            })
    # Also check inside groups
    for shape in slide.shapes:
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                if hasattr(sub, 'shape_type'):
                    if sub.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        text = ""
                        if sub.has_text_frame:
                            text = sub.text_frame.text.strip()
                        auto_shapes.append({
                            'name': sub.name,
                            'text': text,
                            'left': sub.left,
                            'top': sub.top,
                            'width': sub.width,
                            'height': sub.height,
                            'shape': sub,
                        })
                    elif sub.shape_type == MSO_SHAPE_TYPE.LINE or sub.shape_type == 9:
                        connectors.append({
                            'name': sub.name,
                            'left': sub.left,
                            'top': sub.top,
                            'width': sub.width,
                            'height': sub.height,
                            'element': sub._element,
                        })

    print(f"INFO: Found {len(auto_shapes)} auto shapes and {len(connectors)} connectors on slide 2")

    # -------------------------------------------------------
    # Component 1: Three tier rectangles with correct labels (0.30 points)
    # Expected tier labels: 'Presentation Layer', 'Business Logic Layer', 'Data Layer'
    # Tier rectangles should be larger/wider than component boxes
    # -------------------------------------------------------
    try:
        tier_labels = ['Presentation Layer', 'Business Logic Layer', 'Data Layer']
        found_tiers = {}

        for ashape in auto_shapes:
            text_lower = ashape['text'].lower()
            for label in tier_labels:
                if label.lower() in text_lower:
                    found_tiers[label] = ashape

        tier_count = len(found_tiers)
        print(f"  Tier rectangles found: {tier_count}/3 — {list(found_tiers.keys())}")

        if tier_count == 3:
            # Verify vertical ordering: Presentation < Business Logic < Data (by top position)
            tops = [found_tiers[l]['top'] for l in tier_labels]
            if tops[0] < tops[1] < tops[2]:
                print(f"PASS: Component 1 — All 3 tier rectangles present with correct vertical order (0.30 pts)")
                total_score += 0.30
            else:
                # Partial: all 3 found but wrong order
                print(f"PARTIAL: Component 1 — 3 tiers found but vertical order incorrect: tops={tops}")
                total_score += 0.15
        elif tier_count >= 1:
            # Partial credit for partial tier finding
            partial = round(0.10 * tier_count, 2)
            print(f"PARTIAL: Component 1 — Only {tier_count}/3 tier rectangles found ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — No tier rectangles found on slide 2")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # -------------------------------------------------------
    # Component 2: Six component boxes with correct labels (0.40 points)
    # Expected: 'Web App', 'Mobile App' (in Presentation Layer)
    #           'API Gateway', 'Microservices' (in Business Logic Layer)
    #           'PostgreSQL', 'Redis Cache' (in Data Layer)
    # -------------------------------------------------------
    try:
        component_labels = ['Web App', 'Mobile App', 'API Gateway', 'Microservices', 'PostgreSQL', 'Redis Cache']
        found_components = {}

        for ashape in auto_shapes:
            text_lower = ashape['text'].lower()
            for label in component_labels:
                if label.lower() in text_lower:
                    found_components[label] = ashape

        comp_count = len(found_components)
        print(f"  Component boxes found: {comp_count}/6 — {list(found_components.keys())}")

        if comp_count == 6:
            # Verify that components are positioned within their respective tiers
            # Check that Web App and Mobile App are above API Gateway/Microservices,
            # which are above PostgreSQL/Redis Cache
            top_tier_comps = [found_components.get('Web App'), found_components.get('Mobile App')]
            mid_tier_comps = [found_components.get('API Gateway'), found_components.get('Microservices')]
            bot_tier_comps = [found_components.get('PostgreSQL'), found_components.get('Redis Cache')]

            if all(top_tier_comps) and all(mid_tier_comps) and all(bot_tier_comps):
                top_max = max(c['top'] for c in top_tier_comps)
                mid_min = min(c['top'] for c in mid_tier_comps)
                mid_max = max(c['top'] for c in mid_tier_comps)
                bot_min = min(c['top'] for c in bot_tier_comps)

                if top_max < mid_min and mid_max < bot_min:
                    print(f"PASS: Component 2 — All 6 component boxes with correct tier placement (0.40 pts)")
                    total_score += 0.40
                else:
                    print(f"PARTIAL: Component 2 — All 6 boxes found but tier placement incorrect (0.30 pts)")
                    total_score += 0.30
            else:
                print(f"PASS: Component 2 — All 6 component boxes present (0.40 pts)")
                total_score += 0.40
        elif comp_count >= 1:
            partial = round(0.40 * comp_count / 6, 2)
            print(f"PARTIAL: Component 2 — Only {comp_count}/6 component boxes found ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No component boxes found on slide 2")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # -------------------------------------------------------
    # Component 3: Two downward arrow connectors between tiers (0.30 points)
    # Must be line/connector shapes with downward direction (height > 0)
    # and should have arrow heads (tailEnd with triangle type)
    # -------------------------------------------------------
    try:
        arrow_connectors = []
        for conn in connectors:
            # A downward connector has height > 0
            # Check for arrow head in XML
            el = conn['element']
            ln = el.find('.//' + qn('a:ln'))
            has_arrow = _check_arrow_head(ln)

            # Accept connectors that are roughly vertical (width much less than height or width == 0)
            is_vertical = (conn['width'] == 0) or (conn['height'] > 0 and abs(conn['width']) < conn['height'])

            if is_vertical and conn['height'] > 0:
                arrow_connectors.append({
                    'has_arrow_head': has_arrow,
                    'top': conn['top'],
                    'height': conn['height'],
                    'name': conn['name'],
                })

        print(f"  Vertical connectors found: {len(arrow_connectors)}")
        for ac in arrow_connectors:
            print(f"    {ac['name']}: has_arrow_head={ac['has_arrow_head']}, top={ac['top']}, h={ac['height']}")

        # Count connectors with proper arrow heads
        arrows_with_heads = [a for a in arrow_connectors if a['has_arrow_head']]

        if len(arrow_connectors) >= 2:
            if len(arrows_with_heads) >= 2:
                print(f"PASS: Component 3 — {len(arrows_with_heads)} downward arrow connectors found (0.30 pts)")
                total_score += 0.30
            else:
                # Connectors exist but missing arrow heads
                print(f"PARTIAL: Component 3 — {len(arrow_connectors)} vertical connectors but only {len(arrows_with_heads)} have arrow heads (0.20 pts)")
                total_score += 0.20
        elif len(arrow_connectors) == 1:
            if arrows_with_heads:
                print(f"PARTIAL: Component 3 — Only 1 downward arrow connector found (0.15 pts)")
                total_score += 0.15
            else:
                print(f"PARTIAL: Component 3 — 1 vertical connector without arrow head (0.10 pts)")
                total_score += 0.10
        else:
            print(f"FAIL: Component 3 — No vertical connectors found on slide 2")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entrypoint
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/Tech_Architecture.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
