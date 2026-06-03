"""
Reward Script: Machine Learning Lecture Series (15 slides)
Task ID: impress_wf_069
Domain: libreoffice_impress
Scoring:
  Component 1 (0.15): Slide count == 15
  Component 2 (0.10): Slide 1 title contains "Machine Learning Fundamentals"
  Component 3 (0.10): Slide 3 has 3 oval shapes (Venn diagram: AI, ML, DL)
  Component 4 (0.10): Slide 5 has scatter dots (ovals >= 5) and lines
  Component 5 (0.10): Slide 6 has tree structure (ovals + connector lines)
  Component 6 (0.15): Slide 7 has neural network (ovals >= 8 and lines >= 10)
  Component 7 (0.10): Slide 10 has a table for evaluation metrics
  Component 8 (0.10): Slide 11 has confusion matrix (>= 4 rectangle-like shapes)
  Component 9 (0.10): Blue #1565C0 used in fills or backgrounds
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_069'
FILE_NAME = 'ML_Course.pptx'


def count_shape_types(slide):
    """Count ovals, lines, rectangles, tables, and textboxes on a slide."""
    ovals = 0
    lines = 0
    rects = 0
    tables = 0
    textboxes = 0
    for s in slide.shapes:
        try:
            name = s.name
            st = s.shape_type
            if st == 19:  # TABLE
                tables += 1
            elif 'Oval' in name:
                ovals += 1
            elif 'Connector' in name or st == 9:  # LINE
                lines += 1
            elif 'Rectangle' in name or 'Rounded' in name:
                rects += 1
            elif st == 17:  # TEXT_BOX
                textboxes += 1
        except Exception:
            pass
    return {'ovals': ovals, 'lines': lines, 'rects': rects, 'tables': tables, 'textboxes': textboxes}


def get_all_text(slide):
    """Get all text content from a slide."""
    texts = []
    for s in slide.shapes:
        try:
            if s.has_text_frame and s.text_frame.text.strip():
                texts.append(s.text_frame.text.strip())
        except Exception:
            pass
    return texts


def check_blue_usage(prs):
    """Check if blue #1565C0 is used in shape fills or slide backgrounds."""
    target = '1565C0'
    found_count = 0
    for slide in prs.slides:
        # Check slide background
        try:
            fill = slide.background.fill
            if fill.type == 1:
                if str(fill.fore_color.rgb) == target:
                    found_count += 1
        except Exception:
            pass
        # Check shape fills
        for s in slide.shapes:
            try:
                if hasattr(s, 'fill') and s.fill.type == 1:
                    if str(s.fill.fore_color.rgb) == target:
                        found_count += 1
            except Exception:
                pass
    return found_count


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Slide count == 15 (0.15 points)
    try:
        if num_slides == 15:
            print(f"PASS: Component 1 -- Slide count is 15 (0.15 pts)")
            total_score += 0.15
        elif num_slides >= 12:
            partial = 0.08
            print(f"PARTIAL: Component 1 -- Slide count is {num_slides}, expected 15 ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- Slide count is {num_slides}, expected 15")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Gate: Need at least 3 slides for remaining checks
    if num_slides < 3:
        print(f"GATE: Only {num_slides} slides, cannot check slide-specific components")
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: Slide 1 title contains "Machine Learning Fundamentals" (0.10 points)
    try:
        slide1_texts = get_all_text(prs.slides[0])
        slide1_all_text = ' '.join(slide1_texts).lower()
        if 'machine learning fundamentals' in slide1_all_text:
            print(f"PASS: Component 2 -- Slide 1 title 'Machine Learning Fundamentals' found (0.10 pts)")
            total_score += 0.10
        elif 'machine learning' in slide1_all_text:
            print(f"PARTIAL: Component 2 -- Slide 1 has 'machine learning' but not full title (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 2 -- Slide 1 text: {slide1_texts[:3]}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Slide 3 has 3 oval shapes (Venn diagram: AI, ML, DL) (0.10 points)
    try:
        if num_slides >= 3:
            counts3 = count_shape_types(prs.slides[2])
            texts3 = get_all_text(prs.slides[2])
            texts3_lower = ' '.join(texts3).lower()
            has_3_ovals = counts3['ovals'] >= 3
            has_labels = all(label in texts3_lower for label in ['ai', 'ml', 'dl'])
            if has_3_ovals and has_labels:
                print(f"PASS: Component 3 -- Slide 3 has {counts3['ovals']} ovals with AI/ML/DL labels (0.10 pts)")
                total_score += 0.10
            elif has_3_ovals:
                print(f"PARTIAL: Component 3 -- Slide 3 has {counts3['ovals']} ovals but missing labels (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 3 -- Slide 3 has {counts3['ovals']} ovals (need >= 3)")
        else:
            print(f"FAIL: Component 3 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Slide 5 has scatter dots (ovals >= 5) and lines (0.10 points)
    try:
        if num_slides >= 5:
            counts5 = count_shape_types(prs.slides[4])
            has_ovals = counts5['ovals'] >= 5
            has_lines = counts5['lines'] >= 1
            if has_ovals and has_lines:
                print(f"PASS: Component 4 -- Slide 5 has {counts5['ovals']} ovals and {counts5['lines']} lines (0.10 pts)")
                total_score += 0.10
            elif has_ovals or has_lines:
                print(f"PARTIAL: Component 4 -- Slide 5 ovals={counts5['ovals']}, lines={counts5['lines']} (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 4 -- Slide 5 ovals={counts5['ovals']}, lines={counts5['lines']}")
        else:
            print(f"FAIL: Component 4 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Slide 6 has tree structure (ovals + connector lines) (0.10 points)
    try:
        if num_slides >= 6:
            counts6 = count_shape_types(prs.slides[5])
            has_ovals = counts6['ovals'] >= 2
            has_lines = counts6['lines'] >= 2
            has_rects = counts6['rects'] >= 2
            if has_ovals and has_lines and has_rects:
                print(f"PASS: Component 5 -- Slide 6 tree: ovals={counts6['ovals']}, lines={counts6['lines']}, rects={counts6['rects']} (0.10 pts)")
                total_score += 0.10
            elif (has_ovals or has_rects) and has_lines:
                print(f"PARTIAL: Component 5 -- Slide 6 partial tree structure (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 5 -- Slide 6 ovals={counts6['ovals']}, lines={counts6['lines']}, rects={counts6['rects']}")
        else:
            print(f"FAIL: Component 5 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # Component 6: Slide 7 has neural network (ovals >= 8, lines >= 10) (0.15 points)
    try:
        if num_slides >= 7:
            counts7 = count_shape_types(prs.slides[6])
            has_ovals = counts7['ovals'] >= 8
            has_lines = counts7['lines'] >= 10
            if has_ovals and has_lines:
                print(f"PASS: Component 6 -- Slide 7 neural net: ovals={counts7['ovals']}, lines={counts7['lines']} (0.15 pts)")
                total_score += 0.15
            elif counts7['ovals'] >= 4 and counts7['lines'] >= 4:
                print(f"PARTIAL: Component 6 -- Slide 7 partial neural net: ovals={counts7['ovals']}, lines={counts7['lines']} (0.08 pts)")
                total_score += 0.08
            else:
                print(f"FAIL: Component 6 -- Slide 7 ovals={counts7['ovals']}, lines={counts7['lines']}")
        else:
            print(f"FAIL: Component 6 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    # Component 7: Slide 10 has a table (evaluation metrics) (0.10 points)
    try:
        if num_slides >= 10:
            counts10 = count_shape_types(prs.slides[9])
            if counts10['tables'] >= 1:
                print(f"PASS: Component 7 -- Slide 10 has {counts10['tables']} table(s) (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 7 -- Slide 10 has no tables")
        else:
            print(f"FAIL: Component 7 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 7 -- {e}")

    # Component 8: Slide 11 has confusion matrix (>= 4 colored rect-like shapes beyond header) (0.10 points)
    try:
        if num_slides >= 11:
            counts11 = count_shape_types(prs.slides[10])
            texts11 = get_all_text(prs.slides[10])
            texts11_lower = ' '.join(texts11).lower()
            # Confusion matrix: need at least 4 rectangle shapes (the 2x2 grid) + labels
            has_rects = counts11['rects'] >= 4
            has_matrix_text = 'confusion' in texts11_lower or ('predicted' in texts11_lower and 'actual' in texts11_lower)
            if has_rects and has_matrix_text:
                print(f"PASS: Component 8 -- Slide 11 confusion matrix: rects={counts11['rects']}, labels found (0.10 pts)")
                total_score += 0.10
            elif has_rects:
                print(f"PARTIAL: Component 8 -- Slide 11 has {counts11['rects']} rects but missing labels (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 8 -- Slide 11 rects={counts11['rects']}, text: {texts11_lower[:100]}")
        else:
            print(f"FAIL: Component 8 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 8 -- {e}")

    # Component 9: Blue #1565C0 used in fills or backgrounds (0.10 points)
    try:
        blue_count = check_blue_usage(prs)
        if blue_count >= 3:
            print(f"PASS: Component 9 -- Blue #1565C0 found in {blue_count} locations (0.10 pts)")
            total_score += 0.10
        elif blue_count >= 1:
            print(f"PARTIAL: Component 9 -- Blue #1565C0 found in {blue_count} locations (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 9 -- Blue #1565C0 not found in any fills or backgrounds")
    except Exception as e:
        print(f"ERROR: Component 9 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/Desktop/{FILE_NAME}'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
