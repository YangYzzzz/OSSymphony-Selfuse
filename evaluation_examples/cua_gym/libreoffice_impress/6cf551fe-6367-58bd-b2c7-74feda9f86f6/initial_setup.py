"""
Initial Setup: Create a 9-slide TechStack Pitch presentation.
Slide 5 'Platform Architecture' has title only, empty content area.
Task ID: impress_sales_087
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_087'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.paragraphs[0].text = point
        else:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
    return slide


def add_title_only_slide(prs, title_text):
    """Add a slide with title only and empty content."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    # Add a title text box at the top
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(prs, "TechStack Pitch", "Scalable Cloud-Native Platform Architecture\nQ2 2025 Strategy Review")

    # Slide 2: Company Overview
    add_content_slide(prs, "Company Overview", [
        "Founded in 2019 by former AWS and Google Cloud engineers",
        "Headquartered in San Francisco with offices in Berlin and Singapore",
        "Series C funded at $120M valuation (led by Sequoia Capital)",
        "85 full-time employees across engineering, product, and sales",
        "Serving 340+ enterprise clients in financial services and healthcare",
    ])

    # Slide 3: Market Analysis
    add_content_slide(prs, "Market Analysis", [
        "Total addressable market: $47.2B by 2027 (Gartner, 2024)",
        "Cloud infrastructure spending grew 28% YoY in 2024",
        "Key competitors: Datadog, New Relic, Splunk - avg. churn 12%",
        "Our differentiation: unified observability + automated remediation",
        "Target segment: mid-market SaaS companies (200-2000 employees)",
    ])

    # Slide 4: Product Features
    add_content_slide(prs, "Product Features", [
        "Real-time analytics dashboard with sub-second query latency",
        "AI-powered anomaly detection across 150+ metrics",
        "Automated workflow orchestration with 40+ pre-built integrations",
        "Multi-cloud deployment support (AWS, GCP, Azure)",
        "SOC 2 Type II and HIPAA compliance built-in",
    ])

    # Slide 5: Platform Architecture (TITLE ONLY - task target)
    add_title_only_slide(prs, "Platform Architecture")

    # Slide 6: Development Roadmap
    add_content_slide(prs, "Development Roadmap", [
        "Q2 2025: Launch v3.0 with enhanced ML pipeline",
        "Q3 2025: Kubernetes-native operator and Helm charts",
        "Q4 2025: Edge computing support for IoT workloads",
        "Q1 2026: GraphQL API gateway and developer SDK",
        "Q2 2026: Self-hosted enterprise edition with air-gap support",
    ])

    # Slide 7: Team Structure
    add_content_slide(prs, "Team Structure", [
        "Engineering (42): Platform, Data, ML, DevOps, QA teams",
        "Product (12): Product managers, designers, UX researchers",
        "Sales & Marketing (18): Enterprise sales, demand gen, partnerships",
        "Operations (8): Finance, HR, legal, office management",
        "Advisory Board: 5 industry veterans from AWS, Stripe, Databricks",
    ])

    # Slide 8: Financial Projections
    add_content_slide(prs, "Financial Projections", [
        "2024 ARR: $18.4M (142% growth from $7.6M in 2023)",
        "2025 projected ARR: $38M with 65% gross margins",
        "Customer acquisition cost (CAC): $24K, LTV: $186K (7.75x ratio)",
        "Net revenue retention: 135% driven by platform expansion",
        "Path to profitability: Q3 2026 at $52M ARR run rate",
    ])

    # Slide 9: Contact & Next Steps
    add_content_slide(prs, "Contact & Next Steps", [
        "Schedule a personalized demo at techstack.io/demo",
        "Contact: partnerships@techstack.io | (415) 555-0192",
        "Free 30-day enterprise trial with dedicated onboarding",
        "Reference calls available with Acme Financial and MedTech Corp",
        "Investor deck and technical whitepaper available on request",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
