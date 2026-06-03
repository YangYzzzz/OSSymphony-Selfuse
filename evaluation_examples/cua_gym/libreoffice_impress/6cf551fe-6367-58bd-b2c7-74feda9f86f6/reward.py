"""
Reward Script: Technology Stack Architecture Diagram on Slide 5
Task ID: impress_sales_087
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Three layer rectangles with correct labels
  Component 2 (0.15): UI layer sub-boxes (Web App, Mobile App, API)
  Component 3 (0.15): Business Logic sub-boxes (Analytics Engine, Workflow Manager, Integration Hub)
  Component 4 (0.15): Infrastructure sub-boxes (AWS, Database, CDN)
  Component 5 (0.15): Layer colors match spec (#4CAF50, #2196F3, #FF9800)
  Component 6 (0.15): Vertical arrow shapes connecting layers
"""

import os

from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_087'


def persist_app_state(domain):
    """Send Ctrl+S to persist any unsaved GUI edits."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print(f"PERSIST: ctrl+s sent for {domain}")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def get_shape_text_lower(shape):
    """Get shape text in lowercase, stripped."""
    if hasattr(shape, 'text') and shape.text:
        return shape.text.strip().lower()
    return ""


def get_fill_rgb(shape):
    """Get the fill RGB as uppercase hex string, or None."""
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb).upper()
    except Exception:
        pass
    return None


def is_rounded_rectangle(shape):
    """Check if shape is a rounded rectangle auto shape."""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if hasattr(shape, 'auto_shape_type'):
                return shape.auto_shape_type is not None and shape.auto_shape_type == 5  # ROUNDED_RECTANGLE
    except Exception:
        pass
    return False


def is_arrow_shape(shape):
    """Check if shape is a down arrow or similar arrow auto shape."""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if hasattr(shape, 'auto_shape_type'):
                ast = shape.auto_shape_type
                # DOWN_ARROW = 36, UP_ARROW = 68, other arrow types
                if ast is not None and ast in (36, 37, 38, 39, 68, 69, 70, 71):
                    return True
            # Also check by name pattern as fallback
            name = (shape.name or "").lower()
            if "arrow" in name:
                return True
    except Exception:
        pass
    return False


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # Slide 5 (0-indexed)

    # Collect all shapes on slide 5 by type
    rounded_rects = []
    arrows = []
    for shape in slide.shapes:
        if is_rounded_rectangle(shape):
            rounded_rects.append(shape)
        elif is_arrow_shape(shape):
            arrows.append(shape)

    # Build text-to-shape map for rounded rectangles
    rect_by_text = {}
    for r in rounded_rects:
        txt = get_shape_text_lower(r)
        if txt:
            rect_by_text[txt] = r

    print(f"INFO: Found {len(rounded_rects)} rounded rectangles, {len(arrows)} arrows on slide 5")
    print(f"INFO: Rounded rect texts: {list(rect_by_text.keys())}")

    # Component 1: Three layer rectangles with correct labels (0.25 points)
    # These are the main layer boxes: "User Interface", "Business Logic", "Infrastructure"
    try:
        layer_labels = ["user interface", "business logic", "infrastructure"]
        found_layers = []
        for label in layer_labels:
            if label in rect_by_text:
                found_layers.append(label)

        if len(found_layers) == 3:
            print(f"PASS: Component 1 -- All 3 layer rectangles found: {found_layers} (0.25 pts)")
            total_score += 0.25
        elif len(found_layers) >= 1:
            partial = round(0.25 * len(found_layers) / 3, 2)
            print(f"PARTIAL: Component 1 -- {len(found_layers)}/3 layer rectangles found: {found_layers} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- No layer rectangles found (expected: {layer_labels})")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: UI layer sub-boxes (0.15 points)
    # Expected: "Web App", "Mobile App", "API"
    try:
        ui_subs = ["web app", "mobile app", "api"]
        found_ui = [s for s in ui_subs if s in rect_by_text]
        if len(found_ui) == 3:
            print(f"PASS: Component 2 -- All UI sub-boxes found: {found_ui} (0.15 pts)")
            total_score += 0.15
        elif len(found_ui) >= 1:
            partial = round(0.15 * len(found_ui) / 3, 2)
            print(f"PARTIAL: Component 2 -- {len(found_ui)}/3 UI sub-boxes: {found_ui} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 -- No UI sub-boxes found (expected: {ui_subs})")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Business Logic sub-boxes (0.15 points)
    # Expected: "Analytics Engine", "Workflow Manager", "Integration Hub"
    try:
        bl_subs = ["analytics engine", "workflow manager", "integration hub"]
        found_bl = [s for s in bl_subs if s in rect_by_text]
        if len(found_bl) == 3:
            print(f"PASS: Component 3 -- All Business Logic sub-boxes found: {found_bl} (0.15 pts)")
            total_score += 0.15
        elif len(found_bl) >= 1:
            partial = round(0.15 * len(found_bl) / 3, 2)
            print(f"PARTIAL: Component 3 -- {len(found_bl)}/3 BL sub-boxes: {found_bl} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 -- No Business Logic sub-boxes found (expected: {bl_subs})")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Infrastructure sub-boxes (0.15 points)
    # Expected: "AWS", "Database", "CDN"
    try:
        infra_subs = ["aws", "database", "cdn"]
        found_infra = [s for s in infra_subs if s in rect_by_text]
        if len(found_infra) == 3:
            print(f"PASS: Component 4 -- All Infrastructure sub-boxes found: {found_infra} (0.15 pts)")
            total_score += 0.15
        elif len(found_infra) >= 1:
            partial = round(0.15 * len(found_infra) / 3, 2)
            print(f"PARTIAL: Component 4 -- {len(found_infra)}/3 Infra sub-boxes: {found_infra} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 -- No Infrastructure sub-boxes found (expected: {infra_subs})")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Layer colors match spec (0.15 points)
    # UI=#4CAF50, BL=#2196F3, Infra=#FF9800
    try:
        color_checks = {
            "user interface": "4CAF50",
            "business logic": "2196F3",
            "infrastructure": "FF9800",
        }
        colors_correct = 0
        for label, expected_color in color_checks.items():
            if label in rect_by_text:
                actual_rgb = get_fill_rgb(rect_by_text[label])
                if actual_rgb and actual_rgb == expected_color.upper():
                    colors_correct += 1
                    print(f"  Color OK: '{label}' has fill {actual_rgb}")
                else:
                    print(f"  Color MISMATCH: '{label}' expected {expected_color}, got {actual_rgb}")
            else:
                print(f"  Color SKIP: '{label}' shape not found")

        if colors_correct == 3:
            print(f"PASS: Component 5 -- All 3 layer colors correct (0.15 pts)")
            total_score += 0.15
        elif colors_correct >= 1:
            partial = round(0.15 * colors_correct / 3, 2)
            print(f"PARTIAL: Component 5 -- {colors_correct}/3 colors correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 -- No layer colors matched spec")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # Component 6: Vertical arrow shapes connecting layers (0.15 points)
    # Task requires vertical arrows between the 3 layers — expect at least 2 arrows
    try:
        num_arrows = len(arrows)
        if num_arrows >= 4:
            print(f"PASS: Component 6 -- {num_arrows} arrow shapes found (0.15 pts)")
            total_score += 0.15
        elif num_arrows >= 2:
            partial = round(0.15 * 0.7, 2)  # 70% credit for 2-3 arrows
            print(f"PARTIAL: Component 6 -- {num_arrows} arrows found (need >=4 for full) ({partial} pts)")
            total_score += partial
        elif num_arrows >= 1:
            partial = round(0.15 * 0.3, 2)
            print(f"PARTIAL: Component 6 -- Only {num_arrows} arrow(s) found ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 6 -- No arrow shapes found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist any unsaved state before verification
persist_app_state("libreoffice_impress")

# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
