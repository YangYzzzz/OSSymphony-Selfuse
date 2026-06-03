"""
Initial Setup: Brand Launch presentation with company logo on slide 1
Task ID: impress_gf3_020
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_020'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
LOGO_PATH = f'{WORKDIR}/company_logo.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_logo():
    """Create a simple company logo image."""
    img = Image.new('RGBA', (400, 240), (255, 255, 255, 0))
    draw = ImageDraw.Draw(img)
    # Draw a rounded rectangle background
    draw.rounded_rectangle([(10, 10), (390, 230)], radius=20, fill=(0, 102, 179))
    # Draw company initials
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 72)
        small_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
    except OSError:
        font = ImageFont.load_default()
        small_font = font
    draw.text((100, 50), "NV", fill=(255, 255, 255), font=font)
    draw.text((90, 150), "NovaVision", fill=(200, 225, 255), font=small_font)
    img.save(LOGO_PATH)
    print(f'Logo created: {LOGO_PATH}')


def create_initial():
    create_logo()

    prs = Presentation()
    # Standard 16:9 slide dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title with Logo ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    bg1 = slide1.background.fill
    bg1.solid()
    bg1.fore_color.rgb = RGBColor(0x0A, 0x0A, 0x2E)

    # Add company logo - 5cm x 3cm centered in upper portion
    logo_w = Cm(5)
    logo_h = Cm(3)
    logo_left = (prs.slide_width - logo_w) // 2
    logo_top = Cm(3)
    slide1.shapes.add_picture(LOGO_PATH, logo_left, logo_top, logo_w, logo_h)

    # Title text below logo
    txBox = slide1.shapes.add_textbox(Inches(2), Cm(8), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "NovaVision Brand Launch 2026"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    p2 = tf.add_paragraph()
    p2.text = "Redefining Digital Experiences"
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.name = "Arial"
    run2.font.size = Pt(20)
    run2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xFF)

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    bg2 = slide2.background.fill
    bg2.solid()
    bg2.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    title2 = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
    tf2 = title2.text_frame
    p = tf2.paragraphs[0]
    p.text = "Agenda"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x0A, 0x0A, 0x2E)

    agenda_box = slide2.shapes.add_textbox(Inches(1.5), Inches(2), Inches(9), Inches(4))
    atf = agenda_box.text_frame
    atf.word_wrap = True
    items = [
        "1. Brand Identity & Vision",
        "2. Market Positioning Strategy",
        "3. Product Lineup Overview",
        "4. Go-to-Market Timeline",
        "5. Q&A and Next Steps"
    ]
    for i, item in enumerate(items):
        if i == 0:
            pa = atf.paragraphs[0]
        else:
            pa = atf.add_paragraph()
        pa.text = item
        pa.space_after = Pt(14)
        r = pa.runs[0]
        r.font.name = "Arial"
        r.font.size = Pt(22)
        r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 3: Brand Identity ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    bg3 = slide3.background.fill
    bg3.solid()
    bg3.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

    title3 = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
    tf3 = title3.text_frame
    p = tf3.paragraphs[0]
    p.text = "Brand Identity & Vision"
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x0A, 0x0A, 0x2E)

    content3 = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(5.5), Inches(4.5))
    ctf = content3.text_frame
    ctf.word_wrap = True
    paragraphs = [
        ("Mission Statement", True),
        ("To empower businesses and individuals with intuitive, AI-driven digital solutions that transform how people interact with technology.", False),
        ("Core Values", True),
        ("Innovation - Pushing boundaries with cutting-edge research", False),
        ("Accessibility - Technology that works for everyone", False),
        ("Sustainability - Building for the long term", False),
    ]
    for i, (text, is_heading) in enumerate(paragraphs):
        if i == 0:
            pa = ctf.paragraphs[0]
        else:
            pa = ctf.add_paragraph()
        pa.text = text
        r = pa.runs[0]
        r.font.name = "Arial"
        if is_heading:
            r.font.size = Pt(20)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0x00, 0x66, 0xB3)
            pa.space_before = Pt(16)
        else:
            r.font.size = Pt(16)
            r.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # --- Slide 4: Market Positioning ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    title4 = slide4.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
    p = title4.text_frame.paragraphs[0]
    p.text = "Market Positioning"
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x0A, 0x0A, 0x2E)

    # Add a table for competitive analysis
    table_shape = slide4.shapes.add_table(5, 4, Inches(1), Inches(2), Inches(11), Inches(4))
    table = table_shape.table
    headers = ["Feature", "NovaVision", "Competitor A", "Competitor B"]
    data = [
        ["AI Integration", "Full Suite", "Limited", "Basic"],
        ["User Experience", "Intuitive", "Complex", "Moderate"],
        ["Price Point", "$49/mo", "$79/mo", "$39/mo"],
        ["Support", "24/7 Premium", "Business Hours", "Email Only"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        from pptx.oxml.ns import qn
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '0066B3'})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)

    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val
            for run in table.cell(r, c).text_frame.paragraphs[0].runs:
                run.font.size = Pt(13)

    # --- Slide 5: Product Timeline ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    bg5 = slide5.background.fill
    bg5.solid()
    bg5.fore_color.rgb = RGBColor(0x0A, 0x0A, 0x2E)

    title5 = slide5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
    p = title5.text_frame.paragraphs[0]
    p.text = "Go-to-Market Timeline"
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    milestones = [
        ("Q1 2026", "Beta Launch & Early Adopter Program"),
        ("Q2 2026", "Public Release & Marketing Campaign"),
        ("Q3 2026", "Enterprise Tier & API Platform"),
        ("Q4 2026", "International Expansion & Partnerships"),
    ]
    for i, (quarter, desc) in enumerate(milestones):
        left = Inches(1.5 + i * 2.8)
        top = Inches(3)
        box = slide5.shapes.add_textbox(left, top, Inches(2.5), Inches(2))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = quarter
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0x00, 0xBB, 0xFF)
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.runs[0]
        r2.font.size = Pt(14)
        r2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    # --- Slide 6: Thank You ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    bg6 = slide6.background.fill
    bg6.solid()
    bg6.fore_color.rgb = RGBColor(0x0A, 0x0A, 0x2E)

    thanks = slide6.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = thanks.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    p2 = tf.add_paragraph()
    p2.text = "Questions? Contact us at hello@novavision.io"
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.runs[0]
    r2.font.name = "Arial"
    r2.font.size = Pt(18)
    r2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xFF)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
