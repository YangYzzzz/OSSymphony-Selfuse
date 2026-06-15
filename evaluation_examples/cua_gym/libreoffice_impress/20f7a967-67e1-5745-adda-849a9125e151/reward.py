"""
Reward Script: Eisenhower Matrix on Slide 6
Task ID: impress_stu_085
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25) — Four colored quadrant rectangles on slide 6
  Component 2 (0.25) — Quadrant title labels: Schedule, Do First, Eliminate, Delegate
  Component 3 (0.20) — Axis labels: Urgency, Importance, High, Low
  Component 4 (0.15) — Each quadrant contains 2 example task bullet points (8 total)
  Component 5 (0.15) — Grid axes: horizontal and vertical divider lines
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_085'


def get_all_text_from_slide(slide):
    """Extract all text content from a slide, including grouped shapes."""
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return texts


def get_colored_rectangles(slide):
    """Find auto-shapes with solid fill on the slide, return list of (name, color_hex, width, height)."""
    rects = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                fill = shape.fill
                if fill.type is not None and fill.type == 1:  # SOLID fill
                    color_hex = str(fill.fore_color.rgb)
                    rects.append({
                        'name': shape.name,
                        'color': color_hex,
                        'width': shape.width,
                        'height': shape.height,
                        'left': shape.left,
                        'top': shape.top,
                    })
            except Exception:
                pass
    return rects


def get_text_boxes_with_content(slide):
    """Get all text boxes (non-placeholder) with their text content."""
    boxes = []
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame:
            all_text = []
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    all_text.append(t)
            if all_text:
                boxes.append({
                    'name': shape.name,
                    'texts': all_text,
                    'left': shape.left,
                    'top': shape.top,
                })
    return boxes


def verify_task(file_path):
    """
    Verify Eisenhower matrix creation on slide 6 with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[5]  # Slide 6 (0-indexed)

    # Component 1: Four colored quadrant rectangles (0.25 points)
    # The matrix should have 4 large filled rectangles representing quadrants.
    # We filter out thin divider lines by requiring minimum size.
    try:
        all_rects = get_colored_rectangles(slide)
        # Quadrant rectangles are large (width > 1 inch = 914400 EMU, height > 1 inch)
        quadrant_rects = [r for r in all_rects if r['width'] > 914400 and r['height'] > 914400]
        num_quadrants = len(quadrant_rects)

        if num_quadrants >= 4:
            # Check that they have distinct colors (at least 3 different)
            colors = set(r['color'] for r in quadrant_rects[:4])
            if len(colors) >= 3:
                print(f"PASS: Component 1 -- Found {num_quadrants} quadrant rectangles with {len(colors)} distinct colors: {colors} (0.25 pts)")
                total_score += 0.25
            else:
                print(f"PARTIAL: Component 1 -- Found {num_quadrants} quadrant rects but only {len(colors)} distinct colors (0.10 pts)")
                total_score += 0.10
        elif num_quadrants >= 2:
            print(f"PARTIAL: Component 1 -- Found only {num_quadrants} quadrant rectangles, expected 4 (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 1 -- Found {num_quadrants} quadrant rectangles, expected 4")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Quadrant title labels (0.25 points)
    # Expected: "Schedule", "Do First", "Eliminate", "Delegate"
    try:
        all_slide_text = ' '.join(get_all_text_from_slide(slide)).lower()
        expected_labels = ['schedule', 'do first', 'eliminate', 'delegate']
        found_labels = [label for label in expected_labels if label in all_slide_text]

        if len(found_labels) == 4:
            print(f"PASS: Component 2 -- All 4 quadrant labels found: {found_labels} (0.25 pts)")
            total_score += 0.25
        elif len(found_labels) >= 2:
            partial = round(0.25 * len(found_labels) / 4, 2)
            print(f"PARTIAL: Component 2 -- Found {len(found_labels)}/4 labels: {found_labels} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 -- Found only {len(found_labels)}/4 quadrant labels: {found_labels}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Axis labels (0.20 points)
    # Expected: "Urgency" on X-axis, "Importance" on Y-axis, "High" and "Low" markers
    try:
        axis_terms = ['urgency', 'importance']
        marker_terms = ['high', 'low']

        found_axis = [t for t in axis_terms if t in all_slide_text]
        found_markers = [t for t in marker_terms if t in all_slide_text]

        axis_score = 0.0
        if len(found_axis) == 2:
            axis_score += 0.12
        elif len(found_axis) == 1:
            axis_score += 0.06

        if len(found_markers) == 2:
            axis_score += 0.08
        elif len(found_markers) == 1:
            axis_score += 0.04

        if axis_score > 0:
            print(f"PASS: Component 3 -- Axis labels: {found_axis}, markers: {found_markers} ({axis_score} pts)")
            total_score += axis_score
        else:
            print(f"FAIL: Component 3 -- No axis labels found")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Example tasks as bullet points (0.15 points)
    # Each quadrant should have at least 2 example tasks (8 total minimum)
    try:
        text_boxes = get_text_boxes_with_content(slide)
        # Find quadrant text boxes: those containing a quadrant label + bullet items
        quadrant_labels_lower = {'schedule', 'do first', 'eliminate', 'delegate'}
        quadrant_task_counts = {}

        for box in text_boxes:
            # The quadrant label should be the FIRST non-empty text line in the box
            first_text = box['texts'][0].lower().strip() if box['texts'] else ''
            matched_label = None
            for label in quadrant_labels_lower:
                if first_text == label:
                    matched_label = label
                    break
            if matched_label is None:
                continue
            # Count bullet-like items (all lines after the title label)
            bullet_items = []
            for t in box['texts'][1:]:
                t_stripped = t.strip()
                if len(t_stripped) > 5:
                    bullet_items.append(t_stripped)
            if bullet_items:
                quadrant_task_counts[matched_label] = len(bullet_items)

        total_tasks = sum(quadrant_task_counts.values())
        quadrants_with_2_plus = sum(1 for v in quadrant_task_counts.values() if v >= 2)

        if quadrants_with_2_plus >= 4 and total_tasks >= 8:
            print(f"PASS: Component 4 -- All 4 quadrants have 2+ tasks, total={total_tasks} (0.15 pts)")
            total_score += 0.15
        elif total_tasks >= 4:
            partial = round(0.15 * min(total_tasks, 8) / 8, 2)
            print(f"PARTIAL: Component 4 -- {len(quadrant_task_counts)} quadrants with tasks, total={total_tasks} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 -- Found {total_tasks} task items across {len(quadrant_task_counts)} quadrants")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Grid axes / divider lines (0.15 points)
    # The matrix should have at least one horizontal and one vertical thin line/shape acting as axes.
    # These are thin rectangles (one dimension very small).
    try:
        all_rects = get_colored_rectangles(slide)
        # Divider lines: thin in one dimension (< 0.2 inch = 182880 EMU)
        h_dividers = [r for r in all_rects if r['height'] < 182880 and r['width'] > 914400]
        v_dividers = [r for r in all_rects if r['width'] < 182880 and r['height'] > 914400]

        if len(h_dividers) >= 1 and len(v_dividers) >= 1:
            print(f"PASS: Component 5 -- Found {len(h_dividers)} horizontal + {len(v_dividers)} vertical dividers (0.15 pts)")
            total_score += 0.15
        elif len(h_dividers) >= 1 or len(v_dividers) >= 1:
            print(f"PARTIAL: Component 5 -- Found {len(h_dividers)} horizontal, {len(v_dividers)} vertical dividers (0.07 pts)")
            total_score += 0.07
        else:
            print(f"FAIL: Component 5 -- No grid axis dividers found")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
