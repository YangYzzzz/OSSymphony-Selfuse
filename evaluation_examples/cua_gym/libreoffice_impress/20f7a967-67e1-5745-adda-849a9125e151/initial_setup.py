"""
Initial Setup: Time Management presentation with empty slide 6
Task ID: impress_stu_085
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_085'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_body_slide(prs, title_text, body_lines):
    """Add a slide with title and bullet body content."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
    return slide


def add_title_only_slide(prs, title_text):
    """Add a slide with only a title (no body content)."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Time Management Mastery"
    slide1.placeholders[1].text = "Strategies for Peak Productivity\nQ2 2025 Workshop Series"

    # Slide 2: Why Time Management Matters
    add_title_body_slide(prs, "Why Time Management Matters", [
        "82% of professionals lack a formal time management system",
        "Poor time management costs businesses $7,800 per employee annually",
        "Effective time management reduces stress by up to 40%",
        "Leads to better work-life balance and career advancement",
        "Improves decision-making quality under pressure",
    ])

    # Slide 3: Common Time Wasters
    add_title_body_slide(prs, "Common Time Wasters", [
        "Unstructured meetings without agendas (avg. 31 hrs/month wasted)",
        "Email overload and constant inbox checking",
        "Multitasking — reduces productivity by up to 40%",
        "Social media and digital distractions",
        "Lack of clear priorities leading to reactive work",
        "Perfectionism on low-impact tasks",
    ])

    # Slide 4: Goal Setting Framework
    add_title_body_slide(prs, "Goal Setting Framework", [
        "SMART Goals: Specific, Measurable, Achievable, Relevant, Time-bound",
        "Break annual goals into quarterly milestones",
        "Weekly review sessions to track progress",
        "Use OKRs (Objectives and Key Results) for team alignment",
        "Document goals visibly — written goals are 42% more likely achieved",
    ])

    # Slide 5: The Pomodoro Technique
    add_title_body_slide(prs, "The Pomodoro Technique", [
        "Work in focused 25-minute intervals (\"pomodoros\")",
        "Take a 5-minute break after each pomodoro",
        "After 4 pomodoros, take a 15–30 minute longer break",
        "Track completed pomodoros to measure deep work capacity",
        "Best for tasks requiring sustained concentration",
        "Pair with task batching for maximum efficiency",
    ])

    # Slide 6: Prioritization Matrix — EMPTY (task target)
    add_title_only_slide(prs, "Prioritization Matrix")

    # Slide 7: Weekly Planning Template
    add_title_body_slide(prs, "Weekly Planning Template", [
        "Sunday evening: Review upcoming week and set top 3 priorities",
        "Monday morning: Plan daily tasks using time blocking",
        "Mid-week check-in: Adjust priorities based on progress",
        "Friday afternoon: Reflect on accomplishments and lessons",
        "Use a digital calendar synced across all devices",
        "Block 'focus time' — minimum 2-hour uninterrupted blocks",
    ])

    # Slide 8: Key Takeaways
    add_title_body_slide(prs, "Key Takeaways", [
        "Identify and eliminate your top 3 time wasters this week",
        "Implement the Eisenhower Matrix for daily prioritization",
        "Start with one Pomodoro session per day and scale up",
        "Schedule a weekly planning session every Sunday",
        "Track your time for one week to find hidden inefficiencies",
        "Remember: Time management is a skill — it improves with practice",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
