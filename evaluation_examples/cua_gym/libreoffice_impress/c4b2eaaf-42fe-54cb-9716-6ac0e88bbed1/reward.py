"""
FINAL REWARD SCRIPT - SUCCESS
Task: Insert a footnote after 'hypothesis' in paragraph 3 with text 'Preliminary version.'
Generated: 2025-10-17 18:01:47
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation

# Reward Script for verifying the task:
# "Insert a footnote after 'hypothesis' in paragraph 3 with text 'Preliminary version.'"

FOOTNOTE_MARKER = "\u00b9"            # Unicode superscript 1 used as the footnote marker
EXPECTED_FOOTNOTE_TEXT = "preliminary version."  # Expected footnote body text (case-insensitive)


def verify_footnote_task(file_path: str) -> float:
    """Verify that a PPTX file contains:
    1. A footnote marker (¹) immediately following the word 'hypothesis' in paragraph 3
    2. A footnote body whose text is exactly 'Preliminary version.'

    Returns a progressive score between 0.0 and 1.0.
    """
    print(f"Checking file: {file_path}")

    # ----------- Basic file checks (no points awarded) -----------
    if not os.path.exists(file_path):
        print("✗ File does not exist")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error opening presentation: {e}")
        return 0.0

    # ----------- Verification logic -----------
    marker_after_hypothesis = False  # 0.5 points
    footnote_body_correct   = False  # 0.5 points

    # Iterate over every slide and shape to locate the required elements
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            # Only consider shapes with text frames
            if not getattr(shape, "has_text_frame", False):
                continue
            if not shape.has_text_frame:
                continue

            tf = shape.text_frame
            for para_idx, para in enumerate(tf.paragraphs):
                # Reconstruct full paragraph text from runs (preserves marker)
                paragraph_text = "".join(run.text for run in para.runs)
                lower_text = paragraph_text.lower()

                # Requirement 1: marker (¹) immediately after the word 'hypothesis'
                if "hypothesis" in lower_text:
                    pos = lower_text.find("hypothesis")
                    tail = paragraph_text[pos + len("hypothesis"):]
                    if FOOTNOTE_MARKER in tail:
                        marker_after_hypothesis = True
                        print(
                            f"✓ Found footnote marker after 'hypothesis' in slide {slide_idx + 1}, "
                            f"shape {shape_idx}, paragraph {para_idx}"
                        )
                        print(f"  Paragraph text: {paragraph_text}")

                # Requirement 2: footnote body text 'Preliminary version.' (may include marker)
                if FOOTNOTE_MARKER in paragraph_text:
                    cleaned = paragraph_text.replace(FOOTNOTE_MARKER, "").strip().lower()
                    if EXPECTED_FOOTNOTE_TEXT in cleaned:
                        footnote_body_correct = True
                        print(
                            f"✓ Found footnote body text in slide {slide_idx + 1}, "
                            f"shape {shape_idx}, paragraph {para_idx}: {paragraph_text}"
                        )

    # ----------- Scoring (progressive) -----------
    score = 0.0

    if marker_after_hypothesis:
        score += 0.5
    else:
        print("✗ Footnote marker after 'hypothesis' not found")

    if footnote_body_correct:
        score += 0.5
    else:
        print("✗ Correct footnote body text not found")

    final_score = min(score, 1.0)
    print(f"Final score: {final_score}")
    return final_score


# -------------------- Script Entry Point --------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/insert_a_footnote_after_hypothesis_in_paragraph_3_with_text_preliminary_version.pptx"
    reward = verify_footnote_task(FILE_PATH)
    print(f"REWARD: {reward}")
