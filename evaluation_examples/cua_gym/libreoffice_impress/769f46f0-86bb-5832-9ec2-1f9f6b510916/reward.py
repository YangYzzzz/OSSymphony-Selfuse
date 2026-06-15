"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 4 of my Impress deck, there’s a table labeled “Table 1.” The first row still shows the default headers, but I need them changed to the exact text strings “T1”, “T2”, “T3”, and “T4” (left-to-right). What steps do I follow to swap those labels in LibreOffice Impress?
Generated: 2025-09-10 13:07:34
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
import re
from pptx import Presentation

EXPECTED_HEADERS = ["T1", "T2", "T3", "T4"]


def verify_table_headers_task(file_path: str) -> float:
    """Verify that on slide 4 of the given presentation there is a table
    (preferably named 'Table 1') whose first-row cells are exactly
    the strings T1, T2, T3 and T4 from left to right.

    Returns a progressive score between 0.0 and 1.0 and prints a
    detailed breakdown of the verification steps. The function awards
    points ONLY for actual task achievements – no points for natural
    conditions (e.g. file existence, slide existence beyond what is
    required)."""

    print(f"Starting verification for file: {file_path}")
    total_score = 0.0
    max_score = 1.0

    # --- Requirement weights (must sum to 1.0) ---
    WEIGHTS = {
        "slide": 0.1,      # Slide 4 exists
        "table": 0.1,      # Table present on slide 4 (named 'Table 1' preferred)
        "headers": 0.8     # First-row headers exactly T1..T4 (left→right)
    }

    # 1) Load presentation (NO POINTS – prerequisite only)
    if not os.path.exists(file_path):
        print("✗ File does not exist → task failed")
        print("REWARD: 0.0")
        return 0.0
    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded with {len(prs.slides)} slides (no points)")
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # 2) Verify slide 4 exists (index 3)
    if len(prs.slides) >= 4:
        print("✓ Slide 4 exists (0.1 points)")
        total_score += WEIGHTS["slide"]
        slide4 = prs.slides[3]
    else:
        print("✗ Slide 4 not found → remaining checks impossible")
        print(f"REWARD: {total_score}")
        return total_score

    # 3) Look for a table on slide 4
    table_shape = None
    preferred_found = False
    for shape in slide4.shapes:
        if shape.has_table:
            shape_name = getattr(shape, 'name', '') or ''
            # Prefer table explicitly named 'Table 1'
            if re.search(r"Table\s*1", shape_name, re.IGNORECASE):
                table_shape = shape
                preferred_found = True
                break
            # Fallback: remember first table if preferred not found yet
            if table_shape is None:
                table_shape = shape
    if table_shape is None:
        print("✗ No table found on slide 4")
        print(f"REWARD: {total_score}")
        return total_score

    # Table found → award table points
    if preferred_found:
        print("✓ Table named 'Table 1' found (0.1 points)")
    else:
        print("✓ Table found on slide 4 (0.1 points)")
    total_score += WEIGHTS["table"]

    # 4) Verify first-row headers
    table = table_shape.table
    if len(table.rows) == 0:
        print("✗ Table has no rows → cannot verify headers")
        print(f"REWARD: {total_score}")
        return total_score

    first_row_texts = [cell.text.strip() for cell in table.rows[0].cells]
    print("First-row texts detected:", first_row_texts)

    header_matches = 0
    for idx, expected in enumerate(EXPECTED_HEADERS):
        if idx < len(first_row_texts) and first_row_texts[idx] == expected:
            header_matches += 1
            print(f"  ✓ Cell {idx} matches expected '{expected}'")
        else:
            found = first_row_texts[idx] if idx < len(first_row_texts) else "<missing>"
            print(f"  ✗ Cell {idx} expected '{expected}', found '{found}'")

    # Proportional points for header correctness
    header_score = WEIGHTS["headers"] * (header_matches / len(EXPECTED_HEADERS))
    total_score += header_score
    print(f"Headers matched: {header_matches}/{len(EXPECTED_HEADERS)} → {header_score:.2f} points")

    # Final score (capped at 1.0 and rounded for neatness)
    final_score = round(min(total_score, max_score), 4)
    print(f"Total score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score


# ---------------------------
# Run verification when script executed directly
# ---------------------------
if __name__ == "__main__":
    # Path to the presentation file (adjust if necessary)
    PRESENTATION_PATH = "/home/user/on_slide_4_of_my_impress_deck_theres_a_table_labeled_table_1_the_first_row_still_shows_the_default_h_golden.pptx"
    verify_table_headers_task(PRESENTATION_PATH)
