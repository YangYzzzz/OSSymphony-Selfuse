"""
Initial Setup: Hide footer text and date placeholders on master slide
Task ID: impress_ma_047
Domain: libreoffice_impress

Creates a 12-slide business presentation with all three footer area items
(date, footer text, slide number) visible on the master slide.
"""

import os
import shlex
import subprocess
import time
import shutil
import zipfile
from xml.etree import ElementTree as ET
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_047'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # Slide content data
    slide_data = [
        {"layout": 0, "title": "Q4 2025 Strategic Review",
         "subtitle": "Global Operations & Performance Analysis\nPresented by Sarah Chen, VP of Operations"},
        {"layout": 1, "title": "Executive Summary",
         "body": "Revenue grew 18% YoY to $127.3M\nOperating margin improved to 23.4%\nCustomer base expanded to 45,000+ accounts\nNew product launches drove 30% of growth\nAPAC region exceeded targets by 12%"},
        {"layout": 1, "title": "Revenue Breakdown by Region",
         "body": "North America: $52.1M (41% share)\nEurope: $38.2M (30% share)\nAsia Pacific: $24.7M (19% share)\nLatin America: $8.1M (6% share)\nMiddle East & Africa: $4.2M (4% share)"},
        {"layout": 1, "title": "Product Portfolio Performance",
         "body": "Enterprise Suite: $68.5M (+22% YoY)\nCloud Platform: $34.8M (+45% YoY)\nProfessional Services: $15.2M (+8% YoY)\nSupport & Maintenance: $8.8M (+3% YoY)"},
        {"layout": 1, "title": "Key Customer Wins",
         "body": "Meridian Healthcare \u2014 $4.2M enterprise deal\nTechVista Solutions \u2014 $2.8M cloud migration\nGlobal Dynamics Inc \u2014 $3.1M multi-year contract\nAtlas Financial Group \u2014 $1.9M expansion\nPinnacle Manufacturing \u2014 $2.4M platform upgrade"},
        {"layout": 1, "title": "Operational Metrics",
         "body": "Customer satisfaction: 92% (up from 88%)\nEmployee retention: 94.2%\nAverage deal cycle: 47 days (reduced from 62)\nSupport ticket resolution: 4.2 hours avg\nSystem uptime: 99.97%"},
        {"layout": 1, "title": "Technology Infrastructure",
         "body": "Migrated 78% of workloads to Kubernetes\nReduced deployment time by 65%\nImplemented zero-trust security framework\nLaunched AI-powered analytics dashboard\nCompleted SOC 2 Type II certification"},
        {"layout": 1, "title": "People & Culture",
         "body": "Headcount: 1,247 employees (net +182)\nDiversity: 43% women in leadership roles\nTraining: 32 hours per employee average\nRemote/Hybrid: 67% flexible work arrangements\nEngagement score: 4.3/5.0"},
        {"layout": 1, "title": "Risk Assessment",
         "body": "Supply chain disruptions \u2014 Mitigated via dual sourcing\nRegulatory changes in EMEA \u2014 Legal team monitoring\nTalent competition in AI/ML \u2014 Enhanced compensation packages\nCurrency fluctuation exposure \u2014 Hedging strategy active\nCybersecurity threats \u2014 Increased investment by 40%"},
        {"layout": 1, "title": "Q1 2026 Priorities",
         "body": "Launch next-gen Analytics Platform v3.0\nExpand APAC sales team by 25 headcount\nAchieve ISO 27001 certification\nClose 3 strategic partnership deals\nReduce customer churn below 5%"},
        {"layout": 1, "title": "Financial Outlook",
         "body": "FY2026 Revenue Target: $152M (+19%)\nGross Margin Target: 72%\nR&D Investment: $28M (18% of revenue)\nCapital Expenditure: $12M\nFree Cash Flow Target: $35M"},
        {"layout": 1, "title": "Thank You & Questions",
         "body": "Contact: strategy@globalops.com\nNext Review: April 15, 2026\nMaterials available on SharePoint\n\nPrepared by the Strategy & Operations Team"},
    ]

    for i, sd in enumerate(slide_data):
        layout_idx = sd["layout"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        if slide.shapes.title:
            slide.shapes.title.text = sd["title"]
            for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                run.font.size = Pt(28) if layout_idx == 0 else Pt(24)
                run.font.bold = True

        if layout_idx == 0 and "subtitle" in sd:
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = sd["subtitle"]
        elif "body" in sd:
            if len(slide.placeholders) > 1:
                body_ph = slide.placeholders[1]
                tf = body_ph.text_frame
                tf.clear()
                lines = sd["body"].split("\n")
                for j, line in enumerate(lines):
                    if j == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = line
                    p.level = 0
                    for run in p.runs:
                        run.font.size = Pt(16)

    prs.save(OUTPUT)
    print(f"Base presentation saved with {len(slide_data)} slides")

    # Post-process via ZIP/XML to:
    # 1. Add footer placeholders to slide master
    # 2. Set hf element to make all three visible
    temp_path = OUTPUT + '.tmp'
    shutil.move(OUTPUT, temp_path)

    with zipfile.ZipFile(temp_path, 'r') as zin:
        with zipfile.ZipFile(OUTPUT, 'w') as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)

                if item.filename == 'ppt/presentation.xml':
                    root = ET.fromstring(data)
                    # Register namespaces to preserve them
                    ET.register_namespace('', NS_P)
                    ET.register_namespace('a', NS_A)
                    ET.register_namespace('r', NS_R)

                    # Remove existing hf elements
                    for hf in root.findall(f'{{{NS_P}}}hf'):
                        root.remove(hf)
                    # Add hf with all items visible
                    hf_el = ET.SubElement(root, f'{{{NS_P}}}hf')
                    hf_el.set('sldNum', '0')  # visible
                    hf_el.set('dt', '0')       # visible
                    hf_el.set('ftr', '0')      # visible
                    hf_el.set('hdr', '1')      # header hidden
                    data = ET.tostring(root, xml_declaration=True, encoding='UTF-8')

                elif item.filename.startswith('ppt/slideMasters/slideMaster'):
                    root = ET.fromstring(data)
                    sp_tree = root.find(f'{{{NS_P}}}cSld/{{{NS_A}}}spTree')
                    if sp_tree is not None:
                        # Check existing placeholder types
                        existing = set()
                        for sp in sp_tree.findall(f'{{{NS_P}}}sp'):
                            for ph in sp.findall(f'.//{{{NS_P}}}ph'):
                                existing.add(ph.get('type', ''))

                        slide_w = 9144000  # default 10 inches
                        slide_h = 6858000  # default 7.5 inches
                        footer_y = slide_h - 457200
                        footer_h = 365125
                        third_w = slide_w // 3

                        # Add date placeholder (left)
                        if 'dt' not in existing:
                            dt_sp = ET.fromstring(f'''<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}" xmlns:r="{NS_R}">
                                <p:nvSpPr>
                                    <p:cNvPr id="110" name="Date Placeholder 10"/>
                                    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
                                    <p:nvPr><p:ph type="dt" sz="half" idx="10"/></p:nvPr>
                                </p:nvSpPr>
                                <p:spPr>
                                    <a:xfrm>
                                        <a:off x="0" y="{footer_y}"/>
                                        <a:ext cx="{third_w}" cy="{footer_h}"/>
                                    </a:xfrm>
                                </p:spPr>
                                <p:txBody>
                                    <a:bodyPr/>
                                    <a:lstStyle/>
                                    <a:p><a:r><a:rPr lang="en-US" sz="1000"/><a:t>12/15/2025</a:t></a:r></a:p>
                                </p:txBody>
                            </p:sp>''')
                            sp_tree.append(dt_sp)
                            print("Added date placeholder to master")

                        # Add footer text placeholder (center)
                        if 'ftr' not in existing:
                            ftr_sp = ET.fromstring(f'''<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}" xmlns:r="{NS_R}">
                                <p:nvSpPr>
                                    <p:cNvPr id="111" name="Footer Placeholder 11"/>
                                    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
                                    <p:nvPr><p:ph type="ftr" sz="quarter" idx="11"/></p:nvPr>
                                </p:nvSpPr>
                                <p:spPr>
                                    <a:xfrm>
                                        <a:off x="{third_w}" y="{footer_y}"/>
                                        <a:ext cx="{third_w}" cy="{footer_h}"/>
                                    </a:xfrm>
                                </p:spPr>
                                <p:txBody>
                                    <a:bodyPr/>
                                    <a:lstStyle/>
                                    <a:p><a:r><a:rPr lang="en-US" sz="1000"/><a:t>Global Operations - Confidential</a:t></a:r></a:p>
                                </p:txBody>
                            </p:sp>''')
                            sp_tree.append(ftr_sp)
                            print("Added footer text placeholder to master")

                        # Add slide number placeholder (right)
                        if 'sldNum' not in existing:
                            sn_sp = ET.fromstring(f'''<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}" xmlns:r="{NS_R}">
                                <p:nvSpPr>
                                    <p:cNvPr id="112" name="Slide Number Placeholder 12"/>
                                    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
                                    <p:nvPr><p:ph type="sldNum" sz="quarter" idx="12"/></p:nvPr>
                                </p:nvSpPr>
                                <p:spPr>
                                    <a:xfrm>
                                        <a:off x="{third_w * 2}" y="{footer_y}"/>
                                        <a:ext cx="{third_w}" cy="{footer_h}"/>
                                    </a:xfrm>
                                </p:spPr>
                                <p:txBody>
                                    <a:bodyPr/>
                                    <a:lstStyle/>
                                    <a:p>
                                        <a:fld type="slidenum">
                                            <a:rPr lang="en-US" sz="1000"/>
                                            <a:t>&lt;#&gt;</a:t>
                                        </a:fld>
                                    </a:p>
                                </p:txBody>
                            </p:sp>''')
                            sp_tree.append(sn_sp)
                            print("Added slide number placeholder to master")

                    data = ET.tostring(root, xml_declaration=True, encoding='UTF-8')

                zout.writestr(item, data)

    os.remove(temp_path)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
