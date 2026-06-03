"""
Initial Setup: Export all slides of a presentation as individual PNG images
Task ID: osworld_impress_export_image_003
Domain: libreoffice_impress

Creates a 5-slide business review presentation. The agent must then export
each slide as a separate PNG file (5 PNG files total).
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_export_image_003'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_text(slide, title_text, subtitle_text=None, title_color=None):
    """Helper to set title and optional subtitle on a slide."""
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title_text
        run.font.bold = True
        run.font.size = Pt(32)
        if title_color:
            run.font.color.rgb = title_color

    if subtitle_text and len(slide.placeholders) > 1:
        try:
            ph = slide.placeholders[1]
            ph.text_frame.clear()
            p = ph.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = subtitle_text
            run.font.size = Pt(18)
        except Exception:
            pass


def add_text_box(slide, text, left, top, width, height,
                 font_size=16, bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Add a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    # Standard widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    DARK_BLUE = RGBColor(0x1F, 0x3B, 0x6E)
    ACCENT_BLUE = RGBColor(0x27, 0x6B, 0xBF)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
    DARK_TEXT = RGBColor(0x22, 0x22, 0x22)
    GREEN = RGBColor(0x2E, 0x7D, 0x32)
    RED = RGBColor(0xC6, 0x28, 0x28)

    # ── Slide 1: Title Slide ──────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])

    bg1 = slide1.background.fill
    bg1.solid()
    bg1.fore_color.rgb = DARK_BLUE

    if slide1.shapes.title:
        tf = slide1.shapes.title.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "Q2 2025 Business Review"
        run.font.bold = True
        run.font.size = Pt(40)
        run.font.color.rgb = WHITE

    if len(slide1.placeholders) > 1:
        try:
            ph = slide1.placeholders[1]
            ph.text_frame.clear()
            p = ph.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = "Meridian Technologies · June 30, 2025"
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(0xCC, 0xD6, 0xF4)
        except Exception:
            pass

    # ── Slide 2: Executive Summary ───────────────────────────────────────────
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])

    bg2 = slide2.background.fill
    bg2.solid()
    bg2.fore_color.rgb = WHITE

    if slide2.shapes.title:
        tf = slide2.shapes.title.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Executive Summary"
        run.font.bold = True
        run.font.size = Pt(32)
        run.font.color.rgb = DARK_BLUE

    bullets = [
        "Total Revenue: $4.82M (↑ 12.4% YoY)",
        "Net Profit Margin: 18.7% (↑ 2.1pp)",
        "New Customer Acquisitions: 347",
        "Customer Retention Rate: 91.3%",
        "Operating Expenses: $3.92M (within budget)",
        "Headcount: 214 employees (+18 QoQ)",
    ]
    if len(slide2.placeholders) > 1:
        try:
            ph = slide2.placeholders[1]
            ph.text_frame.clear()
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = ph.text_frame.paragraphs[0]
                else:
                    p = ph.text_frame.add_paragraph()
                p.level = 0
                run = p.add_run()
                run.text = bullet
                run.font.size = Pt(18)
                run.font.color.rgb = DARK_TEXT
        except Exception:
            pass

    # ── Slide 3: Sales Performance ───────────────────────────────────────────
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])

    bg3 = slide3.background.fill
    bg3.solid()
    bg3.fore_color.rgb = LIGHT_GRAY

    if slide3.shapes.title:
        tf = slide3.shapes.title.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Sales Performance by Region"
        run.font.bold = True
        run.font.size = Pt(32)
        run.font.color.rgb = DARK_BLUE

    # Add a table with sales data
    table_shape = slide3.shapes.add_table(
        6, 4,
        Inches(1.0), Inches(1.8),
        Inches(11.0), Inches(4.5)
    )
    table = table_shape.table

    headers = ["Region", "Q2 2025 Revenue", "Q2 2024 Revenue", "YoY Growth"]
    data = [
        ["North America",  "$1,842,500",  "$1,620,000",  "+13.7%"],
        ["Europe",         "$1,234,800",  "$1,110,200",  "+11.2%"],
        ["Asia Pacific",   "$980,400",    "$845,300",    "+16.0%"],
        ["Latin America",  "$562,100",    "$504,700",    "+11.4%"],
        ["Middle East",    "$200,200",    "$176,800",    "+13.2%"],
    ]

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = WHITE
        # Header background via XML
        from pptx.oxml.ns import qn
        from lxml import etree
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', '1F3B6E')

    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(13)
                run.font.color.rgb = DARK_TEXT
                if c == 3:  # Growth column
                    run.font.color.rgb = GREEN

    # ── Slide 4: Product Highlights ──────────────────────────────────────────
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])

    bg4 = slide4.background.fill
    bg4.solid()
    bg4.fore_color.rgb = WHITE

    if slide4.shapes.title:
        tf = slide4.shapes.title.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Product Highlights & Roadmap"
        run.font.bold = True
        run.font.size = Pt(32)
        run.font.color.rgb = DARK_BLUE

    products = [
        ("MeridianCloud Pro 3.2",   "Released Apr 15 · 2,100 active licenses · NPS 72"),
        ("DataBridge Enterprise",   "Released May 28 · 480 enterprise clients · SLA 99.97%"),
        ("SecureVault 2.0",         "Launched Jun 10 · 650 deployments · Zero critical incidents"),
        ("Analytics Studio (Beta)", "Q3 2025 GA target · 190 beta participants · Positive feedback"),
    ]
    if len(slide4.placeholders) > 1:
        try:
            ph = slide4.placeholders[1]
            ph.text_frame.clear()
            for i, (name, desc) in enumerate(products):
                if i == 0:
                    p = ph.text_frame.paragraphs[0]
                else:
                    p = ph.text_frame.add_paragraph()
                p.level = 0
                run = p.add_run()
                run.text = f"{name}: {desc}"
                run.font.size = Pt(16)
                run.font.color.rgb = DARK_TEXT
        except Exception:
            pass

    # ── Slide 5: Outlook & Next Steps ────────────────────────────────────────
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])

    bg5 = slide5.background.fill
    bg5.solid()
    bg5.fore_color.rgb = DARK_BLUE

    if slide5.shapes.title:
        tf = slide5.shapes.title.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Q3 2025 Outlook & Next Steps"
        run.font.bold = True
        run.font.size = Pt(32)
        run.font.color.rgb = WHITE

    next_steps = [
        "Accelerate Asia Pacific expansion — hire 15 regional sales reps by Aug 2025",
        "Launch Analytics Studio GA — complete final QA and certification by Sep 2025",
        "Achieve $5.2M Q3 revenue target (↑ 7.9% QoQ)",
        "Reduce operating expense ratio from 81.3% to 79.5%",
        "Partner with 3 new system integrators in EMEA region",
        "Initiate Series C fundraising process — target close Nov 2025",
    ]
    if len(slide5.placeholders) > 1:
        try:
            ph = slide5.placeholders[1]
            ph.text_frame.clear()
            for i, step in enumerate(next_steps):
                if i == 0:
                    p = ph.text_frame.paragraphs[0]
                else:
                    p = ph.text_frame.add_paragraph()
                p.level = 0
                run = p.add_run()
                run.text = step
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(0xCC, 0xD6, 0xF4)
        except Exception:
            pass

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Slide count: {len(prs.slides)}')

    # GUI-ready startup — open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
