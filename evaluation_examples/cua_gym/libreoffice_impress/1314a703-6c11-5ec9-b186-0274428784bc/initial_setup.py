"""
Initial Setup: Create a 5-slide presentation with title 'WELCOME' on slide 1 (default spacing)
Task ID: impress_tct_079
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_079'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide with "WELCOME" ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "WELCOME"
    for run in slide1.shapes.title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
    slide1.placeholders[1].text = "Annual Strategy Review 2025"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    items = [
        "Company Performance Overview",
        "Market Trends & Competitive Analysis",
        "Product Roadmap for Q3-Q4",
        "Financial Projections",
        "Open Discussion & Q&A",
    ]
    body2.paragraphs[0].text = items[0]
    for item in items[1:]:
        p = body2.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 3: Performance Highlights ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Performance Highlights"
    body3 = slide3.placeholders[1].text_frame
    highlights = [
        "Revenue grew 18% YoY to $124.5M",
        "Customer retention rate improved to 94.2%",
        "Launched 3 new product lines across APAC region",
        "Employee satisfaction score reached 4.6/5.0",
        "Reduced operational costs by 12% through automation",
    ]
    body3.paragraphs[0].text = highlights[0]
    for h in highlights[1:]:
        p = body3.add_paragraph()
        p.text = h
        p.level = 0

    # --- Slide 4: Market Analysis ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Market Analysis"
    body4 = slide4.placeholders[1].text_frame
    analysis = [
        "Total addressable market estimated at $8.2B by 2026",
        "Key competitor acquisitions signal market consolidation",
        "Rising demand for AI-integrated solutions (+32% CAGR)",
        "Regulatory changes in EU creating new compliance needs",
    ]
    body4.paragraphs[0].text = analysis[0]
    for a in analysis[1:]:
        p = body4.add_paragraph()
        p.text = a
        p.level = 0

    # --- Slide 5: Next Steps ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Next Steps"
    body5 = slide5.placeholders[1].text_frame
    steps = [
        "Finalize Q3 budget allocations by June 15",
        "Schedule cross-functional planning workshops",
        "Launch customer feedback survey for new features",
        "Prepare investor update materials for board meeting",
    ]
    body5.paragraphs[0].text = steps[0]
    for s in steps[1:]:
        p = body5.add_paragraph()
        p.text = s
        p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
