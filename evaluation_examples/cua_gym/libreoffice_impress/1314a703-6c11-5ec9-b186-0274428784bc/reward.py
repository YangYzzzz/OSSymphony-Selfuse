"""
Reward Script: Adjust character spacing (kerning) of slide 1 title to expanded by 3pt
Task ID: impress_tct_079
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Title on slide 1 has character spacing attribute set
  Component 2 (0.4): Character spacing value is exactly 300 (3pt in hundredths of a point)
  Component 3 (0.2): All runs in the title have the expanded spacing applied
"""

import os
from pptx import Presentation
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_079'


def persist_app_state(domain: str):
    """Save any unsaved GUI state before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0

    Task: Adjust character spacing of slide 1 title 'WELCOME' to expanded by 3pt.
    In OOXML, character spacing is stored as the 'spc' attribute on <a:rPr> elements,
    in hundredths of a point. 3pt = 300.
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Must have at least 1 slide
    if len(prs.slides) < 1:
        print("FAIL: Presentation has no slides")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]

    # Find the title shape on slide 1
    title_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name and 'title' in shape.name.lower():
            title_shape = shape
            break

    # Fallback: look for shape with "WELCOME" text
    if title_shape is None:
        for shape in slide.shapes:
            if shape.has_text_frame and 'WELCOME' in shape.text_frame.text.upper():
                title_shape = shape
                break

    if title_shape is None:
        print("FAIL: No title shape found on slide 1")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Title text should be WELCOME
    title_text = title_shape.text_frame.text.strip()
    if 'WELCOME' not in title_text.upper():
        print(f"FAIL: Title text is '{title_text}', expected 'WELCOME'")
        print("REWARD: 0.0")
        return 0.0

    # Collect all non-empty runs in the title
    all_runs = []
    for para in title_shape.text_frame.paragraphs:
        for run in para.runs:
            if (run.text or "").strip():
                all_runs.append(run)

    if len(all_runs) == 0:
        print("FAIL: No text runs found in title")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Found title shape '{title_shape.name}' with text '{title_text}' and {len(all_runs)} run(s)")

    # Component 1: At least one run in the title has EXPANDED character spacing (positive spc > 0) (0.4 points)
    # The initial file may have spc=-1 (near-zero condensed from LO save). We need positive expanded spacing.
    try:
        any_has_expanded_spacing = False
        for run in all_runs:
            rPr = run.font._element
            spc_val = rPr.attrib.get('spc', None)
            if spc_val is not None:
                try:
                    spc_int = int(spc_val)
                    if spc_int > 0:
                        any_has_expanded_spacing = True
                    print(f"  Run '{run.text}': spc={spc_val} ({'expanded' if spc_int > 0 else 'condensed/zero'})")
                except ValueError:
                    print(f"  Run '{run.text}': spc={spc_val} (non-integer)")
            else:
                print(f"  Run '{run.text}': spc=NOT SET (default)")

        if any_has_expanded_spacing:
            print(f"PASS: Component 1 -- Expanded character spacing found on title runs (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 -- No expanded character spacing (spc > 0) found on any title run")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Character spacing value is exactly 300 (3pt) on at least one run (0.4 points)
    try:
        any_correct_value = False
        for run in all_runs:
            rPr = run.font._element
            spc_val = rPr.attrib.get('spc', None)
            if spc_val is not None:
                try:
                    spc_int = int(spc_val)
                    if spc_int == 300:
                        any_correct_value = True
                except ValueError:
                    pass

        if any_correct_value:
            print(f"PASS: Component 2 -- Character spacing is exactly 300 (3pt) (0.4 pts)")
            total_score += 0.4
        else:
            # Check if there's any spacing set but wrong value
            found_values = []
            for run in all_runs:
                rPr = run.font._element
                spc_val = rPr.attrib.get('spc', None)
                if spc_val is not None:
                    found_values.append(spc_val)
            if found_values:
                print(f"FAIL: Component 2 -- Character spacing values found: {found_values}, expected 300")
            else:
                print(f"FAIL: Component 2 -- No character spacing values found")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: ALL runs in the title have expanded spacing of 300 (0.2 points)
    try:
        all_runs_correct = True
        for run in all_runs:
            rPr = run.font._element
            spc_val = rPr.attrib.get('spc', None)
            if spc_val is None or int(spc_val) != 300:
                all_runs_correct = False
                break

        if all_runs_correct:
            print(f"PASS: Component 3 -- All {len(all_runs)} title run(s) have spacing=300 (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 3 -- Not all title runs have spacing=300")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist unsaved GUI state before scoring
persist_app_state("libreoffice_impress")

# Run verification
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
