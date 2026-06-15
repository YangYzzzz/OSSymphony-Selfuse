"""
Initial Setup: Create a 15-slide photo slideshow presentation.
Task ID: impress_tm_008
Domain: libreoffice_impress
Slide 6 has default advance (On Mouse Click).
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_008'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slideshow theme data
    slides_data = [
        {
            "title": "Wanderlust: A Photo Journey",
            "subtitle": "Capturing moments from around the world",
            "bg_color": RGBColor(0x1A, 0x1A, 0x2E),
            "title_color": RGBColor(0xFF, 0xFF, 0xFF),
        },
        {
            "title": "Sunrise over Santorini",
            "description": "The golden light breaking through the Aegean mist, casting warm hues across the whitewashed villages perched on volcanic cliffs. Photographed during an early morning hike along the caldera rim in October 2024.",
            "bg_color": RGBColor(0xFF, 0xE0, 0xB2),
            "title_color": RGBColor(0x8B, 0x45, 0x13),
        },
        {
            "title": "Tokyo Neon Nights",
            "description": "Shinjuku's electric energy captured in the reflections of rain-soaked streets. Neon signs in Japanese and English create layers of color that transform the urban landscape into a living canvas of light.",
            "bg_color": RGBColor(0x1A, 0x0A, 0x2E),
            "title_color": RGBColor(0xE0, 0x40, 0xFF),
        },
        {
            "title": "Patagonian Wilderness",
            "description": "Torres del Paine at dawn, with jagged granite peaks piercing the clouds. The turquoise waters of Lake Pehoe mirror the mountains, creating a symmetry that feels almost surreal in its perfection.",
            "bg_color": RGBColor(0x2E, 0x4A, 0x62),
            "title_color": RGBColor(0xB0, 0xE0, 0xE6),
        },
        {
            "title": "Marrakech Medina",
            "description": "The labyrinthine souks alive with color -- saffron, indigo, and vermillion spices piled high in ceramic bowls. Handwoven carpets drape from ancient doorways while the call to prayer echoes overhead.",
            "bg_color": RGBColor(0xC4, 0x5A, 0x12),
            "title_color": RGBColor(0xFF, 0xF8, 0xDC),
        },
        {
            "title": "Northern Lights over Tromso",
            "description": "The aurora borealis dancing in ribbons of emerald green and violet above the frozen Norwegian fjords. A lone cabin sits beneath the cosmic display, its windows glowing amber against the arctic night.",
            "bg_color": RGBColor(0x0A, 0x1A, 0x2A),
            "title_color": RGBColor(0x00, 0xFF, 0x7F),
        },
        {
            "title": "Kyoto Cherry Blossoms",
            "description": "The Philosopher's Path draped in delicate pink sakura. Petals drift lazily onto the canal below, creating a pastel carpet that moves with the gentle spring current. Taken during peak bloom in early April.",
            "bg_color": RGBColor(0xFF, 0xE4, 0xE1),
            "title_color": RGBColor(0x8B, 0x00, 0x45),
        },
        {
            "title": "Sahara Desert Caravan",
            "description": "Silhouettes of camels stretching across endless golden dunes at sunset. The Merzouga erg transforms into waves of amber and bronze as the sun descends, while Berber guides navigate by ancient starlight.",
            "bg_color": RGBColor(0xDA, 0xA5, 0x20),
            "title_color": RGBColor(0x4A, 0x2C, 0x0A),
        },
        {
            "title": "Iceland's Black Sand Beach",
            "description": "Reynisfjara's volcanic shoreline where diamond-like ice chunks wash ashore against jet-black sand. Basalt columns rise like organ pipes from the cliff face, a testament to the island's volcanic origins.",
            "bg_color": RGBColor(0x2A, 0x2A, 0x2A),
            "title_color": RGBColor(0xE0, 0xE0, 0xE0),
        },
        {
            "title": "Venice at Twilight",
            "description": "The Grand Canal reflecting the last light of day as gondolas glide beneath the Rialto Bridge. Renaissance facades glow in shades of terracotta and gold, their foundations kissed by the rising tide.",
            "bg_color": RGBColor(0x2E, 0x1A, 0x47),
            "title_color": RGBColor(0xFF, 0xD7, 0x00),
        },
        {
            "title": "Great Barrier Reef",
            "description": "An underwater kaleidoscope of coral formations and tropical fish. Clownfish dart between anemone tentacles while a sea turtle glides overhead through crystal-clear water illuminated by shafts of sunlight.",
            "bg_color": RGBColor(0x00, 0x6B, 0x8A),
            "title_color": RGBColor(0x7F, 0xFF, 0xD4),
        },
        {
            "title": "Machu Picchu at Dawn",
            "description": "The lost Incan citadel emerging from morning clouds, its precisely cut stone terraces cascading down the mountainside. Llamas graze peacefully among the ancient ruins as mist swirls through the Sacred Valley.",
            "bg_color": RGBColor(0x3A, 0x5A, 0x28),
            "title_color": RGBColor(0xF0, 0xE6, 0x8C),
        },
        {
            "title": "New York City Skyline",
            "description": "Manhattan from the Brooklyn Bridge at blue hour. The Empire State Building and One World Trade Center punctuate the skyline while the East River captures their reflections in rippling patterns of light.",
            "bg_color": RGBColor(0x1C, 0x2B, 0x4A),
            "title_color": RGBColor(0xFF, 0xA5, 0x00),
        },
        {
            "title": "Balinese Rice Terraces",
            "description": "The emerald cascades of Tegallalang carved into the volcanic hillside. Morning light catches water flowing between paddies as farmers in conical hats tend to the ancient subak irrigation system.",
            "bg_color": RGBColor(0x22, 0x8B, 0x22),
            "title_color": RGBColor(0xFF, 0xFF, 0xE0),
        },
        {
            "title": "Thank You",
            "subtitle": "Photos and text by Elena Vasquez | www.wanderlust-gallery.com",
            "bg_color": RGBColor(0x1A, 0x1A, 0x2E),
            "title_color": RGBColor(0xFF, 0xFF, 0xFF),
        },
    ]

    for i, sd in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only

        # Background
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = sd["bg_color"]

        # Title
        title_top = Inches(0.8) if "subtitle" in sd else Inches(0.6)
        add_text_box(slide, Inches(1), title_top, Inches(11), Inches(1.2),
                     sd["title"], font_size=36, bold=True,
                     color=sd["title_color"], alignment=PP_ALIGN.CENTER)

        # Description or subtitle
        if "description" in sd:
            add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(3.5),
                         sd["description"], font_size=20,
                         color=sd["title_color"], alignment=PP_ALIGN.LEFT)

            # Slide number indicator
            add_text_box(slide, Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.5),
                         f"{i + 1} / 15", font_size=12,
                         color=sd["title_color"], alignment=PP_ALIGN.RIGHT)

        if "subtitle" in sd:
            sub_color = RGBColor(0xCC, 0xCC, 0xCC) if i == 0 else RGBColor(0x99, 0x99, 0x99)
            add_text_box(slide, Inches(2), Inches(3.5), Inches(9), Inches(1),
                         sd["subtitle"], font_size=20,
                         color=sub_color, alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
