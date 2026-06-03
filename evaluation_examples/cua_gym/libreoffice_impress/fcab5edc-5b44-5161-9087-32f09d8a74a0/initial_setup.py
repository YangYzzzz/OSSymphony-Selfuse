"""
Initial Setup: Add speaker notes to slides 1-5 of a quarterly update presentation
Task ID: impress_fix_075
Domain: libreoffice_impress

Creates a 10-slide Quarterly Update presentation. Slides 1-5 have NO speaker notes
(the task is to add them). Slides 6-10 have content but also no notes.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_075'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=None):
    """Helper to add a text box to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)
    return txBox


def add_bullet_points(slide, left, top, width, height, items, font_size=16):
    """Add bullet-point text box."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(font_size)
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    add_textbox(slide1, 1.5, 1.5, 7, 1.5, "Quarterly Business Update",
                font_size=36, bold=True, color=(0xFF, 0xFF, 0xFF),
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, 1.5, 3.5, 7, 1, "Q1 2025 Performance Review",
                font_size=22, color=(0xCC, 0xDD, 0xEE),
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, 1.5, 5.0, 7, 0.8, "Presented by Elena Rodriguez, VP of Operations",
                font_size=16, color=(0x99, 0xBB, 0xDD),
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, 1.5, 5.8, 7, 0.5, "March 31, 2025",
                font_size=14, color=(0x88, 0xAA, 0xCC),
                alignment=PP_ALIGN.CENTER)
    # NO notes on slide 1 - that is the task

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide2, 0.5, 0.3, 9, 1, "Agenda",
                font_size=32, bold=True, color=(0x1B, 0x3A, 0x5C))
    add_bullet_points(slide2, 0.8, 1.5, 8.5, 5, [
        "1. Q1 Revenue & Financial Highlights",
        "2. Product Development Milestones",
        "3. Customer Acquisition & Retention Metrics",
        "4. Challenges & Lessons Learned",
        "5. Strategic Roadmap for Q2-Q4",
        "6. Team Updates & Hiring Plans",
        "7. Budget Allocation Review",
        "8. Risk Assessment & Mitigation",
        "9. Key Partnerships Update",
        "10. Q&A Session",
    ], font_size=18)

    # --- Slide 3: Q1 Highlights ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide3, 0.5, 0.3, 9, 1, "Q1 Financial Highlights",
                font_size=32, bold=True, color=(0x1B, 0x3A, 0x5C))
    add_bullet_points(slide3, 0.8, 1.5, 8.5, 5, [
        "Total Revenue: $4.82M (+18% YoY)",
        "Gross Margin: 72.3% (up from 68.1%)",
        "New Enterprise Contracts: 14 signed",
        "Monthly Recurring Revenue: $1.61M",
        "Customer Lifetime Value increased to $38,400",
        "Operating Expenses reduced by 6.2%",
    ], font_size=18)

    # --- Slide 4: Challenges ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide4, 0.5, 0.3, 9, 1, "Challenges Faced in Q1",
                font_size=32, bold=True, color=(0x1B, 0x3A, 0x5C))
    add_bullet_points(slide4, 0.8, 1.5, 8.5, 5, [
        "Supply chain delays impacted hardware delivery by 3 weeks",
        "Engineering team attrition rate rose to 12%",
        "Cloud infrastructure costs exceeded budget by $145K",
        "Two major client implementations ran behind schedule",
        "Compliance audit required unplanned resource allocation",
        "Competitor launched similar feature set in February",
    ], font_size=18)

    # --- Slide 5: Roadmap ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide5, 0.5, 0.3, 9, 1, "Strategic Roadmap: Q2-Q4 2025",
                font_size=32, bold=True, color=(0x1B, 0x3A, 0x5C))
    add_bullet_points(slide5, 0.8, 1.5, 8.5, 5, [
        "Q2: Launch v3.0 platform with AI-powered analytics",
        "Q2: Expand into European market (UK, Germany, France)",
        "Q3: Achieve SOC 2 Type II certification",
        "Q3: Hire 25 additional engineers across 3 teams",
        "Q4: Target $7.5M quarterly revenue milestone",
        "Q4: Establish strategic partnership with Salesforce",
    ], font_size=18)

    # --- Slide 6: Team Updates ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide6, 0.5, 0.3, 9, 1, "Team Updates",
                font_size=32, bold=True, color=(0x1B, 0x3A, 0x5C))
    add_bullet_points(slide6, 0.8, 1.5, 8.5, 5, [
        "Engineering: 48 members, 5 new hires in Q1",
        "Sales: Expanded APAC team with 3 regional leads",
        "Marketing: Rebranded campaign launched in February",
        "Customer Success: NPS improved from 62 to 71",
        "HR: New mentorship program rolled out company-wide",
    ], font_size=18)

    # --- Slide 7: Budget Review ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide7, 0.5, 0.3, 9, 1, "Budget Allocation Review",
                font_size=32, bold=True, color=(0x1B, 0x3A, 0x5C))
    add_bullet_points(slide7, 0.8, 1.5, 8.5, 5, [
        "R&D: $2.1M (43% of total budget)",
        "Sales & Marketing: $1.4M (29%)",
        "Operations: $850K (17%)",
        "General & Administrative: $540K (11%)",
        "Contingency fund remaining: $320K",
    ], font_size=18)

    # --- Slide 8: Risk Assessment ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide8, 0.5, 0.3, 9, 1, "Risk Assessment & Mitigation",
                font_size=32, bold=True, color=(0x1B, 0x3A, 0x5C))
    add_bullet_points(slide8, 0.8, 1.5, 8.5, 5, [
        "HIGH: Talent retention - implementing equity refresh program",
        "MEDIUM: Cloud cost overrun - migrating to reserved instances",
        "MEDIUM: Regulatory changes in EU - legal team monitoring",
        "LOW: Currency fluctuation - hedging strategy in place",
    ], font_size=18)

    # --- Slide 9: Partnerships ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide9, 0.5, 0.3, 9, 1, "Key Partnerships Update",
                font_size=32, bold=True, color=(0x1B, 0x3A, 0x5C))
    add_bullet_points(slide9, 0.8, 1.5, 8.5, 5, [
        "AWS: Advanced Technology Partner status achieved",
        "Deloitte: Co-selling agreement signed for enterprise deals",
        "Stanford AI Lab: Research collaboration on NLP models",
        "TechStars: Mentorship program for internal innovation",
    ], font_size=18)

    # --- Slide 10: Q&A ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = slide10.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    add_textbox(slide10, 1.5, 2.5, 7, 1.5, "Questions & Discussion",
                font_size=36, bold=True, color=(0xFF, 0xFF, 0xFF),
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide10, 1.5, 4.5, 7, 1, "Thank you for your time and engagement",
                font_size=20, color=(0xCC, 0xDD, 0xEE),
                alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
