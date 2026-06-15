"""
Reward Script: Add speaker notes to slides 1-5 of a presentation
Task ID: impress_fix_075
Domain: libreoffice_impress
Scoring: 5 components (0.2 each) — one per slide's speaker notes text
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_075'

# Expected notes for slides 1-5 (0-indexed keys 0-4)
EXPECTED_NOTES = {
    0: 'Welcome the audience, introduce yourself',
    1: 'Explain the agenda',
    2: 'Cover Q1 highlights',
    3: 'Discuss challenges faced',
    4: 'Present the roadmap',
}

POINTS_PER_SLIDE = 0.2


def persist_app_state(domain: str):
    """Save any unsaved GUI edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)

    for slide_idx, expected_text in EXPECTED_NOTES.items():
        slide_num = slide_idx + 1
        # Component: Slide N speaker notes (0.2 points)
        try:
            slide = slides[slide_idx]
            try:
                actual_notes = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                actual_notes = ""

            if actual_notes == expected_text:
                print(f"PASS: Slide {slide_num} notes match exactly ({POINTS_PER_SLIDE} pts)")
                total_score += POINTS_PER_SLIDE
            elif expected_text.lower() in actual_notes.lower():
                # Partial: text is present but may have extra content
                partial = POINTS_PER_SLIDE * 0.5
                print(f"PARTIAL: Slide {slide_num} notes contain expected text but not exact match ({partial} pts)")
                print(f"  Expected: {repr(expected_text)}")
                print(f"  Actual:   {repr(actual_notes)}")
                total_score += partial
            else:
                print(f"FAIL: Slide {slide_num} notes do not match")
                print(f"  Expected: {repr(expected_text)}")
                print(f"  Actual:   {repr(actual_notes)}")
        except Exception as e:
            print(f"ERROR: Slide {slide_num} check failed: {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
