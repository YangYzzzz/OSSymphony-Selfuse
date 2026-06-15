"""
Reward Script: Pricing comparison table with three card shapes on slide 7
Task ID: impress_rp_025
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): Basic card — white body, blue header #3498DB, 'Basic' text, '$9.99/mo'
  Component 2 (0.30): Pro card — #EBF5FB body, gold border #F1C40F ~3pt, blue header, 'Pro' text, '$24.99/mo', featured (taller)
  Component 3 (0.20): Enterprise card — white body, dark header #2C3E50, 'Enterprise' text, 'Contact Us'
  Component 4 (0.15): Three cards arranged horizontally (left positions increasing)
  Component 5 (0.15): All header texts are white (#FFFFFF)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_025'


def color_match(actual_rgb, expected_hex, tolerance=10):
    """Check if an RGBColor matches expected hex within tolerance per channel."""
    if actual_rgb is None:
        return False
    actual_str = str(actual_rgb).upper()
    expected_hex = expected_hex.upper().lstrip('#')
    try:
        ar, ag, ab = int(actual_str[0:2], 16), int(actual_str[2:4], 16), int(actual_str[4:6], 16)
        er, eg, eb = int(expected_hex[0:2], 16), int(expected_hex[2:4], 16), int(expected_hex[4:6], 16)
        return abs(ar - er) <= tolerance and abs(ag - eg) <= tolerance and abs(ab - eb) <= tolerance
    except (ValueError, IndexError):
        return False


def get_shape_fill_color(shape):
    """Get solid fill color of a shape, returns RGBColor or None."""
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # SOLID
            return fill.fore_color.rgb
    except Exception:
        pass
    return None


def get_shape_line_color_and_width(shape):
    """Get line (border) color and width of a shape."""
    try:
        line = shape.line
        if line.fill.type is not None:
            return line.color.rgb, line.width
    except Exception:
        pass
    return None, None


def get_all_text(shape):
    """Extract all text from a shape's text frame."""
    if not shape.has_text_frame:
        return ""
    return " ".join(para.text for para in shape.text_frame.paragraphs).strip()


def get_text_color(shape):
    """Get the font color of the first non-empty run."""
    if not shape.has_text_frame:
        return None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if (run.text or "").strip():
                try:
                    if run.font.color.type is not None:
                        return run.font.color.rgb
                except Exception:
                    pass
    return None


def find_card_groups(slide):
    """
    Identify three pricing card groups on slide 7.
    Each card consists of:
      - A body rectangle (larger, with fill)
      - A header rectangle (smaller, overlapping top of body, with colored fill and text)
      - A text box with price/CTA text
    Returns dict keyed by card name ('Basic', 'Pro', 'Enterprise') with shape info.
    """
    cards = {}
    auto_shapes = []
    text_boxes = []

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            auto_shapes.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            text = get_all_text(shape)
            # Skip the title "Pricing Plans"
            if text and text != "Pricing Plans":
                text_boxes.append(shape)

    # Identify header shapes (auto_shapes with tier name text)
    header_shapes = {}
    body_shapes = []
    for s in auto_shapes:
        text = get_all_text(s)
        if text in ('Basic', 'Pro', 'Enterprise'):
            header_shapes[text] = s
        else:
            body_shapes.append(s)

    # For each identified header, find its body (same left, larger height) and price text
    for name, header in header_shapes.items():
        card = {'header': header, 'body': None, 'price_box': None}

        # Find body: same or very close left position, but taller
        for bs in body_shapes:
            if abs(bs.left - header.left) < Inches(0.2) and bs.height > header.height:
                card['body'] = bs
                break

        # Find price text box: horizontally overlapping with header, below header
        header_center_x = header.left + header.width // 2
        for tb in text_boxes:
            tb_center_x = tb.left + tb.width // 2
            if abs(tb_center_x - header_center_x) < Inches(1.5) and tb.top > header.top:
                card['price_box'] = tb
                break

        cards[name] = card

    return cards


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Check slide count
    if len(prs.slides) < 7:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 7")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[6]  # slide 7 (0-indexed)
    cards = find_card_groups(slide)

    print(f"Found cards: {list(cards.keys())}")
    for name, card in cards.items():
        print(f"  {name}: body={'YES' if card['body'] else 'NO'}, header={'YES' if card['header'] else 'NO'}, price={'YES' if card['price_box'] else 'NO'}")

    # Component 1: Basic card (0.20 points)
    # White body, blue #3498DB header, 'Basic' text, '$9.99/mo' price
    try:
        if 'Basic' in cards and cards['Basic']['body'] and cards['Basic']['price_box']:
            card = cards['Basic']
            sub = 0.0

            # Body fill is white
            body_color = get_shape_fill_color(card['body'])
            if body_color and color_match(body_color, 'FFFFFF'):
                sub += 0.05
                print(f"PASS: Basic body is white ({body_color})")
            else:
                print(f"FAIL: Basic body fill expected white, got {body_color}")

            # Header fill is blue #3498DB
            header_color = get_shape_fill_color(card['header'])
            if header_color and color_match(header_color, '3498DB'):
                sub += 0.05
                print(f"PASS: Basic header is blue ({header_color})")
            else:
                print(f"FAIL: Basic header fill expected #3498DB, got {header_color}")

            # Header text says 'Basic'
            header_text = get_all_text(card['header'])
            if 'Basic' in header_text:
                sub += 0.05
                print(f"PASS: Basic header text found: {repr(header_text)}")
            else:
                print(f"FAIL: Basic header text expected 'Basic', got {repr(header_text)}")

            # Price text says '$9.99/mo'
            price_text = get_all_text(card['price_box'])
            if '$9.99' in price_text and 'mo' in price_text.lower():
                sub += 0.05
                print(f"PASS: Basic price text found: {repr(price_text)}")
            else:
                print(f"FAIL: Basic price expected '$9.99/mo', got {repr(price_text)}")

            total_score += sub
            print(f"Component 1 (Basic card): {sub}/0.20 pts")
        else:
            print(f"FAIL: Component 1 — Basic card not found or incomplete")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Pro card (0.30 points)
    # #EBF5FB body, gold border #F1C40F ~3pt, blue header, 'Pro' text, '$24.99/mo', featured (taller)
    try:
        if 'Pro' in cards and cards['Pro']['body'] and cards['Pro']['price_box']:
            card = cards['Pro']
            sub = 0.0

            # Body fill is light blue #EBF5FB
            body_color = get_shape_fill_color(card['body'])
            if body_color and color_match(body_color, 'EBF5FB'):
                sub += 0.05
                print(f"PASS: Pro body is light blue ({body_color})")
            else:
                print(f"FAIL: Pro body fill expected #EBF5FB, got {body_color}")

            # Gold border #F1C40F with ~3pt width (38100 EMU = 3pt)
            line_color, line_width = get_shape_line_color_and_width(card['body'])
            if line_color and color_match(line_color, 'F1C40F'):
                sub += 0.05
                print(f"PASS: Pro border color is gold ({line_color})")
            else:
                print(f"FAIL: Pro border color expected #F1C40F, got {line_color}")

            if line_width and abs(line_width - 38100) <= 6350:  # ~3pt +/- 0.5pt tolerance
                sub += 0.05
                print(f"PASS: Pro border width ~3pt ({line_width} EMU)")
            else:
                print(f"FAIL: Pro border width expected ~38100 EMU (3pt), got {line_width}")

            # Header fill is blue #3498DB
            header_color = get_shape_fill_color(card['header'])
            if header_color and color_match(header_color, '3498DB'):
                sub += 0.025
                print(f"PASS: Pro header is blue ({header_color})")
            else:
                print(f"FAIL: Pro header fill expected #3498DB, got {header_color}")

            # Header text says 'Pro'
            header_text = get_all_text(card['header'])
            if 'Pro' in header_text:
                sub += 0.025
                print(f"PASS: Pro header text found: {repr(header_text)}")
            else:
                print(f"FAIL: Pro header text expected 'Pro', got {repr(header_text)}")

            # Price text says '$24.99/mo'
            price_text = get_all_text(card['price_box'])
            if '$24.99' in price_text and 'mo' in price_text.lower():
                sub += 0.05
                print(f"PASS: Pro price text found: {repr(price_text)}")
            else:
                print(f"FAIL: Pro price expected '$24.99/mo', got {repr(price_text)}")

            # Featured: Pro body is taller than Basic or Enterprise body
            basic_body = cards.get('Basic', {}).get('body')
            enterprise_body = cards.get('Enterprise', {}).get('body')
            pro_body = card['body']
            if pro_body and (basic_body or enterprise_body):
                ref_height = basic_body.height if basic_body else enterprise_body.height
                if pro_body.height > ref_height:
                    sub += 0.05
                    print(f"PASS: Pro card is taller ({pro_body.height}) than reference ({ref_height})")
                else:
                    print(f"FAIL: Pro card height ({pro_body.height}) not taller than reference ({ref_height})")
            else:
                print(f"FAIL: Cannot compare Pro height — missing reference cards")

            total_score += sub
            print(f"Component 2 (Pro card): {sub}/0.30 pts")
        else:
            print(f"FAIL: Component 2 — Pro card not found or incomplete")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Enterprise card (0.20 points)
    # White body, dark header #2C3E50, 'Enterprise' text, 'Contact Us'
    try:
        if 'Enterprise' in cards and cards['Enterprise']['body'] and cards['Enterprise']['price_box']:
            card = cards['Enterprise']
            sub = 0.0

            # Body fill is white
            body_color = get_shape_fill_color(card['body'])
            if body_color and color_match(body_color, 'FFFFFF'):
                sub += 0.05
                print(f"PASS: Enterprise body is white ({body_color})")
            else:
                print(f"FAIL: Enterprise body fill expected white, got {body_color}")

            # Header fill is dark #2C3E50
            header_color = get_shape_fill_color(card['header'])
            if header_color and color_match(header_color, '2C3E50'):
                sub += 0.05
                print(f"PASS: Enterprise header is dark ({header_color})")
            else:
                print(f"FAIL: Enterprise header fill expected #2C3E50, got {header_color}")

            # Header text says 'Enterprise'
            header_text = get_all_text(card['header'])
            if 'Enterprise' in header_text:
                sub += 0.05
                print(f"PASS: Enterprise header text found: {repr(header_text)}")
            else:
                print(f"FAIL: Enterprise header text expected 'Enterprise', got {repr(header_text)}")

            # Price text says 'Contact Us'
            price_text = get_all_text(card['price_box'])
            if 'contact' in price_text.lower() and 'us' in price_text.lower():
                sub += 0.05
                print(f"PASS: Enterprise CTA text found: {repr(price_text)}")
            else:
                print(f"FAIL: Enterprise CTA expected 'Contact Us', got {repr(price_text)}")

            total_score += sub
            print(f"Component 3 (Enterprise card): {sub}/0.20 pts")
        else:
            print(f"FAIL: Component 3 — Enterprise card not found or incomplete")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Horizontal arrangement (0.15 points)
    # Three cards side by side: left positions should be increasing
    try:
        if len(cards) == 3 and all(cards[n].get('body') for n in ('Basic', 'Pro', 'Enterprise')):
            positions = []
            for name in ('Basic', 'Pro', 'Enterprise'):
                positions.append((name, cards[name]['body'].left))
            positions.sort(key=lambda x: x[1])

            # Check that positions are distinct and increasing
            names_order = [p[0] for p in positions]
            lefts = [p[1] for p in positions]
            if lefts[0] < lefts[1] < lefts[2]:
                # Check reasonable spacing (not overlapping)
                gap1 = lefts[1] - (lefts[0] + cards[positions[0][0]]['body'].width)
                gap2 = lefts[2] - (lefts[1] + cards[positions[1][0]]['body'].width)
                if gap1 > 0 and gap2 > 0:
                    print(f"PASS: Cards arranged horizontally L-to-R: {names_order}, gaps: {gap1}, {gap2}")
                    total_score += 0.15
                else:
                    print(f"FAIL: Cards overlap. Gaps: {gap1}, {gap2}")
            else:
                print(f"FAIL: Cards not in horizontal order. Lefts: {lefts}")
        else:
            print(f"FAIL: Component 4 — need exactly 3 cards with bodies")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: All header texts are white (0.15 points)
    try:
        white_count = 0
        checked = 0
        for name in ('Basic', 'Pro', 'Enterprise'):
            if name in cards and cards[name].get('header'):
                header = cards[name]['header']
                text_color = get_text_color(header)
                checked += 1
                if text_color and color_match(text_color, 'FFFFFF'):
                    white_count += 1
                    print(f"PASS: {name} header text is white ({text_color})")
                else:
                    print(f"FAIL: {name} header text color expected white, got {text_color}")
        if checked == 3 and white_count == 3:
            total_score += 0.15
            print(f"Component 5 (white headers): 0.15/0.15 pts")
        elif checked > 0 and white_count > 0:
            partial = 0.15 * (white_count / 3)
            total_score += partial
            print(f"Component 5 (white headers): {partial:.3f}/0.15 pts (partial)")
        else:
            print(f"FAIL: Component 5 — no white header texts found")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
