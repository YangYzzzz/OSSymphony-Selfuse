"""
Initial Setup: SaaS Pitch Deck with 12 slides, slide 7 has title only
Task ID: impress_rp_025
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_025'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_title_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
    return slide


def add_title_only_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs, "CloudFlow SaaS Platform", "Streamline Your Business Operations\nQ2 2025 Investor Pitch")

    # Slide 2: About Us
    add_title_content_slide(prs, "About CloudFlow", [
        "Founded in 2021 by former Salesforce and AWS engineers",
        "Headquartered in San Francisco with offices in London and Singapore",
        "Over 2,500 enterprise customers across 40 countries",
        "Series C funded — $85M raised to date",
        "Named in Gartner's Magic Quadrant for Cloud Management, 2024",
    ])

    # Slide 3: The Problem
    add_title_content_slide(prs, "The Problem We Solve", [
        "Enterprises waste 30% of cloud spend on unused or misallocated resources",
        "Manual workflow orchestration leads to 12+ hours/week lost per team",
        "Legacy integration tools require 6-9 months for deployment",
        "Data silos across departments reduce decision-making speed by 40%",
    ])

    # Slide 4: Our Solution
    add_title_content_slide(prs, "Our Solution", [
        "AI-powered resource optimization reduces cloud costs by up to 35%",
        "No-code workflow builder deploys in under 2 weeks",
        "Unified data layer connects 150+ enterprise applications",
        "Real-time analytics dashboard with predictive insights",
        "SOC 2 Type II and ISO 27001 certified",
    ])

    # Slide 5: Key Features
    add_title_content_slide(prs, "Key Features", [
        "Smart Resource Allocator — ML-driven capacity planning",
        "FlowBuilder — Drag-and-drop workflow automation",
        "DataBridge — Universal API connector with pre-built templates",
        "InsightHub — Customizable dashboards and scheduled reports",
        "TeamSync — Cross-departmental collaboration workspace",
    ])

    # Slide 6: Customer Success
    add_title_content_slide(prs, "Customer Success Stories", [
        "Meridian Healthcare: Reduced cloud costs by $1.2M annually",
        "Pinnacle Financial: Automated 85% of compliance workflows",
        "Atlas Logistics: Cut integration deployment time from 8 months to 3 weeks",
        "NovaTech Manufacturing: Improved cross-team visibility by 60%",
    ])

    # Slide 7: Pricing Plans — TITLE ONLY, no cards
    add_title_only_slide(prs, "Pricing Plans")

    # Slide 8: Market Opportunity
    add_title_content_slide(prs, "Market Opportunity", [
        "Total Addressable Market: $78B by 2027 (Forrester Research)",
        "Cloud management segment growing at 24% CAGR",
        "Only 15% of mid-market enterprises have adopted AI-driven optimization",
        "Regulatory tailwinds: GDPR, SOX compliance driving automation demand",
    ])

    # Slide 9: Go-to-Market Strategy
    add_title_content_slide(prs, "Go-to-Market Strategy", [
        "Product-led growth with 14-day free trial funnel",
        "Enterprise sales team targeting Fortune 1000 accounts",
        "Strategic partnerships with AWS, Azure, and GCP marketplaces",
        "Channel partner program launching Q3 2025",
        "Annual CloudFlow Summit conference (3,000+ attendees in 2024)",
    ])

    # Slide 10: Team
    add_title_content_slide(prs, "Leadership Team", [
        "Elena Vasquez, CEO — Former VP of Product at Salesforce",
        "James Okonkwo, CTO — Ex-Principal Engineer at AWS",
        "Priya Sharma, CFO — Previously at Goldman Sachs Technology",
        "Marcus Chen, VP Sales — Built $200M ARR org at ServiceNow",
        "Sofia Andersson, VP Engineering — 15 years distributed systems",
    ])

    # Slide 11: Financial Projections
    add_title_content_slide(prs, "Financial Projections", [
        "2024 ARR: $42M (up 115% YoY)",
        "2025 Projected ARR: $78M",
        "Gross margin: 82% (industry-leading)",
        "Net revenue retention: 135%",
        "Path to profitability: Q4 2026",
    ])

    # Slide 12: Call to Action
    add_title_slide(prs, "Let's Build the Future Together", "Contact: partnerships@cloudflow.io\nwww.cloudflow.io")

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
