"""
FINAL REWARD SCRIPT - SUCCESS
Task: Insert a math formula at the caret with content 'x^2 + y^2 = z^2'.
Generated: 2025-10-17 16:46:15
Status: success
Model: azure-o3
Total Steps: 4
"""

from pptx import Presentation
import os
import re

# -----------------------------------------------------------------------------
# Reward script for task:
# "Insert a math formula at the caret with content 'x^2 + y^2 = z^2'."
# -----------------------------------------------------------------------------
# Verification strategy:
# 1. Load the provided PPTX file safely (no points for mere existence/loading).
# 2. Extract every text fragment from all slides.
# 3. Search for the exact formula pattern  x^2 + y^2 = z^2  (case-insensitive, 
#    flexible on spaces, allows either "^2" or the superscript character "²").
#    • If the full expression is found – award full credit (1.0).
# 4. If exact expression isn’t found, award partial credit only for partial
#    presence of tokens (x^2, y^2, z^2, x², …) – progressive scoring.
#    • This guarantees no hard-coded success and no points for natural states.
# -----------------------------------------------------------------------------

FILE_PATH = "/home/user/insert_a_math_formula_at_the_caret_with_content_x2_y2_z2.pptx"

# --------------------------- helper functions ---------------------------------

def load_presentation(path):
    """Safely load a PPTX file and report slide count."""
    if not os.path.exists(path):
        print(f"✗ File not found: {path}")
        return None
    try:
        prs = Presentation(path)
        print(f"✓ Loaded presentation with {len(prs.slides)} slide(s)")
        return prs
    except Exception as exc:
        print(f"✗ Error loading PPTX: {exc}")
        return None

def extract_all_text(prs):
    """Return list of individual texts and one big combined string."""
    texts = []
    for s_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = shape.text or ""
                if txt.strip():
                    texts.append(txt)
                    print(f"  [Slide {s_idx}] text: {repr(txt.strip()[:80])}")
    combined = "\n".join(texts)
    print(f"Total extracted text length: {len(combined)} characters")
    return texts, combined

def score_formula_presence(texts, combined):
    """Return a progressive score based on how well the formula is present."""
    # Regex for flexible detection of the exact formula
    pattern = re.compile(
        r"x\s*(?:\^|[²])\s*2?\s*\+\s*y\s*(?:\^|[²])\s*2?\s*=\s*z\s*(?:\^|[²])\s*2?",
        re.IGNORECASE,
    )

    # Check exact expression in any single text block first
    for text in texts:
        if pattern.search(text):
            print(f"✓ Exact formula found in text block: {repr(text.strip()[:80])}")
            return 1.0  # perfect completion

    # If not found, try across combined text (handles line breaks)
    if pattern.search(combined):
        print("✓ Exact formula found across combined text")
        return 1.0

    # ---------------- partial credit pathway ----------------
    tokens = ["x^2", "y^2", "z^2", "x²", "y²", "z²"]
    token_presence = {tok: (tok.lower() in combined.lower()) for tok in tokens}
    tokens_found = sum(token_presence.values())

    print("Partial token presence:")
    for tok, present in token_presence.items():
        print(f"  {tok}: {'✓' if present else '✗'}")

    # Progressive scoring based on how many distinct tokens are present
    if tokens_found >= 6:          # every variant present – very close, but order wrong
        return 0.8
    elif tokens_found >= 3:        # all three squares present, but no exact formula
        return 0.5
    elif tokens_found >= 1:        # at least one  square present
        return 0.3
    else:
        return 0.0                 # nothing indicative found

# --------------------------- main verification --------------------------------

def verify_task(path):
    print(f"Verifying task for file: {path}\n" + "-"*60)
    prs = load_presentation(path)
    if prs is None:
        print("Verification failed: could not load presentation.")
        return 0.0

    texts, combined = extract_all_text(prs)
    reward = score_formula_presence(texts, combined)

    print("-"*60)
    print(f"Final calculated reward: {reward}")
    return reward

# --------------------------- script execution ---------------------------------
if __name__ == "__main__":
    final_reward = verify_task(FILE_PATH)
    print(f"REWARD: {final_reward}")
