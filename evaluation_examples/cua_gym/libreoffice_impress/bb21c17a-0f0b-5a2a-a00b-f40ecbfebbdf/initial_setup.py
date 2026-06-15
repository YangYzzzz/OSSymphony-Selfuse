"""
Initial Setup: Create a 6-slide presentation with slide 4 titled 'Project Timeline' but empty.
Task ID: impress_stu_039
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_039'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_body_text(slide, text, font_size=18):
    """Add text to the content placeholder (index 1) if it exists."""
    if len(slide.placeholders) > 1:
        ph = slide.placeholders[1]
        tf = ph.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        for run in p.runs:
            run.font.size = Pt(font_size)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Student Research Project"
    slide1.placeholders[1].text = "Analyzing the Impact of Remote Learning on Student Engagement"

    # --- Slide 2: Introduction (Title + Content) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Introduction"
    tf2 = slide2.placeholders[1].text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = "This research examines how remote learning affected student engagement across 12 universities during 2023-2025."
    p2b = tf2.add_paragraph()
    p2b.text = "Key areas of focus include attendance patterns, participation rates, and academic performance metrics."
    p2b.space_before = Pt(12)
    p2c = tf2.add_paragraph()
    p2c.text = "The study covers both undergraduate and graduate programs across STEM and humanities disciplines."
    p2c.space_before = Pt(12)

    # --- Slide 3: Methodology (Title + Content) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Methodology"
    tf3 = slide3.placeholders[1].text_frame
    tf3.word_wrap = True
    tf3.paragraphs[0].text = "Mixed-methods approach combining quantitative survey data (n=2,450) with qualitative interviews (n=85)."
    p3b = tf3.add_paragraph()
    p3b.text = "Data collection period: September 2024 through February 2025."
    p3b.space_before = Pt(12)
    p3c = tf3.add_paragraph()
    p3c.text = "Statistical analysis performed using R (v4.3.2) with significance threshold p < 0.05."
    p3c.space_before = Pt(12)

    # --- Slide 4: Project Timeline (Title Only — EMPTY body) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title as a text box at the top
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf4 = txBox.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Project Timeline"
    p4.alignment = PP_ALIGN.CENTER
    run4 = p4.runs[0]
    run4.font.size = Pt(36)
    run4.font.bold = True
    # NO timeline shapes — the agent's task is to create them

    # --- Slide 5: Results (Title + Content) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Results"
    tf5 = slide5.placeholders[1].text_frame
    tf5.word_wrap = True
    tf5.paragraphs[0].text = "Remote learning reduced in-class participation by 34% but increased asynchronous forum engagement by 28%."
    p5b = tf5.add_paragraph()
    p5b.text = "Students in STEM programs showed greater adaptability to remote formats compared to humanities students."
    p5b.space_before = Pt(12)
    p5c = tf5.add_paragraph()
    p5c.text = "Overall GPA remained stable (3.21 vs 3.18), suggesting compensatory learning strategies."
    p5c.space_before = Pt(12)

    # --- Slide 6: Conclusion (Title + Content) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Conclusion"
    tf6 = slide6.placeholders[1].text_frame
    tf6.word_wrap = True
    tf6.paragraphs[0].text = "The findings suggest that hybrid learning models can maintain academic standards while improving accessibility."
    p6b = tf6.add_paragraph()
    p6b.text = "Recommendations include structured asynchronous activities and regular synchronous check-ins."
    p6b.space_before = Pt(12)
    p6c = tf6.add_paragraph()
    p6c.text = "Future research should explore long-term retention and career readiness outcomes."
    p6c.space_before = Pt(12)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
