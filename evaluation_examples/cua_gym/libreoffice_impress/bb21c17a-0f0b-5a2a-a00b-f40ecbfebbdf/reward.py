"""
Reward Script: Create a timeline on slide 4 with milestones
Task ID: impress_stu_039
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): 5 oval/circle shapes on slide 4
  Component 2 (0.2): Horizontal connecting line (rectangle shape) on slide 4
  Component 3 (0.3): All 5 milestone label texts present on slide 4
  Component 4 (0.2): Labels positioned below the circles
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_039'

# Expected milestone labels (normalized lowercase for matching)
EXPECTED_LABELS = [
    "week 1-2",
    "week 3-4",
    "week 5-6",
    "week 7",
    "week 8",
]

EXPECTED_KEYWORDS = [
    "research",
    "data collection",
    "analysis",
    "writing",
    "presentation",
]


def normalize_text(text):
    """Normalize text for comparison: lowercase, collapse whitespace, strip."""
    import re
    text = text.replace('\x0b', ' ').replace('\n', ' ')
    text = re.sub(r'\s+', ' ', text).strip().lower()
    return text


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)

    # Collect shape info from slide 4
    ovals = []
    rectangles = []
    text_boxes = []

    for shape in slide.shapes:
        # Check for ovals (circle milestones)
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                from pptx.enum.shapes import MSO_SHAPE
                if shape.auto_shape_type == MSO_SHAPE.OVAL:
                    ovals.append(shape)
                elif shape.auto_shape_type == MSO_SHAPE.RECTANGLE:
                    rectangles.append(shape)
        except Exception:
            pass

        # Collect text boxes (not placeholders) for label checking
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            if shape.has_text_frame:
                full_text = ""
                for para in shape.text_frame.paragraphs:
                    full_text += para.text + " "
                full_text = normalize_text(full_text)
                if full_text:
                    text_boxes.append({
                        'text': full_text,
                        'top': shape.top,
                        'left': shape.left,
                        'shape': shape,
                    })

    print(f"INFO: Slide 4 has {len(ovals)} ovals, {len(rectangles)} rectangles, {len(text_boxes)} text boxes")

    # Component 1: 5 oval/circle shapes on slide 4 (0.3 points)
    # Initial slide 4 has 0 ovals, golden has 5
    try:
        oval_count = len(ovals)
        if oval_count >= 5:
            print(f"PASS: Component 1 -- Found {oval_count} oval shapes (need 5) (0.3 pts)")
            total_score += 0.3
        elif oval_count >= 3:
            partial = 0.15
            print(f"PARTIAL: Component 1 -- Found {oval_count}/5 oval shapes ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- Found only {oval_count} oval shapes, need 5")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Horizontal connecting line/rectangle on slide 4 (0.2 points)
    # Initial slide 4 has 0 rectangles, golden has 1 (thin horizontal bar)
    try:
        # Look for a rectangle that is much wider than tall (line-like)
        line_found = False
        for rect in rectangles:
            # A connecting line rectangle is much wider than tall
            if rect.width > rect.height * 5:
                line_found = True
                print(f"PASS: Component 2 -- Found horizontal line shape (w={rect.width}, h={rect.height}) (0.2 pts)")
                total_score += 0.2
                break
        if not line_found:
            # Also check for actual line shapes (freeform connector)
            for shape in slide.shapes:
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.LINE or 'line' in shape.name.lower():
                        line_found = True
                        print(f"PASS: Component 2 -- Found line shape '{shape.name}' (0.2 pts)")
                        total_score += 0.2
                        break
                except Exception:
                    pass
            if not line_found:
                # Check for any wide rectangle even if not in our list
                for shape in slide.shapes:
                    try:
                        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.width > shape.height * 5:
                            line_found = True
                            print(f"PASS: Component 2 -- Found wide shape acting as line (0.2 pts)")
                            total_score += 0.2
                            break
                    except Exception:
                        pass
            if not line_found:
                print(f"FAIL: Component 2 -- No horizontal connecting line found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: All 5 milestone labels present (0.3 points)
    # Initial slide 4 has only "Project Timeline", golden has all 5 milestone labels
    try:
        all_text_on_slide = " ".join([tb['text'] for tb in text_boxes])
        labels_found = 0
        for i, (label, keyword) in enumerate(zip(EXPECTED_LABELS, EXPECTED_KEYWORDS)):
            # Check if both the week reference and the keyword are present
            if label in all_text_on_slide and keyword in all_text_on_slide:
                labels_found += 1
            elif label in all_text_on_slide or keyword in all_text_on_slide:
                # Partial: at least one part found
                labels_found += 0.5

        if labels_found >= 5:
            print(f"PASS: Component 3 -- All 5 milestone labels found (0.3 pts)")
            total_score += 0.3
        elif labels_found >= 3:
            partial = round(0.3 * (labels_found / 5), 2)
            print(f"PARTIAL: Component 3 -- {labels_found}/5 labels found ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 -- Only {labels_found}/5 milestone labels found")
            print(f"  All text on slide 4: {all_text_on_slide[:200]}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Labels positioned below circles (0.2 points)
    # Check that text boxes with milestone labels have top > oval tops
    try:
        if len(ovals) >= 3 and labels_found >= 3:
            # Get average oval top position
            avg_oval_top = sum(o.top for o in ovals) / len(ovals)
            avg_oval_bottom = avg_oval_top + (sum(o.height for o in ovals) / len(ovals))

            # Filter text boxes to only milestone labels (exclude "Project Timeline")
            milestone_textboxes = []
            for tb in text_boxes:
                for label in EXPECTED_LABELS:
                    if label in tb['text']:
                        milestone_textboxes.append(tb)
                        break

            if len(milestone_textboxes) >= 3:
                # Check that milestone labels are below the ovals
                labels_below = sum(1 for tb in milestone_textboxes if tb['top'] > avg_oval_top)
                if labels_below >= len(milestone_textboxes) * 0.8:
                    print(f"PASS: Component 4 -- {labels_below}/{len(milestone_textboxes)} labels below circles (0.2 pts)")
                    total_score += 0.2
                else:
                    print(f"FAIL: Component 4 -- Only {labels_below}/{len(milestone_textboxes)} labels below circles")
            else:
                print(f"FAIL: Component 4 -- Not enough milestone text boxes found ({len(milestone_textboxes)})")
        else:
            print(f"FAIL: Component 4 -- Not enough ovals ({len(ovals)}) or labels ({labels_found}) to check positioning")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
