"""
Initial Setup: Create a 5-slide presentation with 'Cumulative User Growth' title on slide 4, no chart.
Task ID: impress_tct_053
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_053'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # ============================================================
    # Slide 1: Title Slide — "User Growth Analysis"
    # ============================================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "User Growth Analysis"
    slide1.placeholders[1].text = "Q1-Q4 2025 Performance Report"

    # ============================================================
    # Slide 2: "Monthly Signups Overview" with a data table
    # ============================================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Monthly Signups Overview"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1E, 0x88, 0xE5)

    # Table with monthly signup data
    months = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun',
              'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
    signups = [1200, 1450, 1680, 2100, 2350, 2800,
               3150, 3600, 4020, 4500, 5100, 5800]

    tbl_shape = slide2.shapes.add_table(
        3, 7, Inches(0.3), Inches(1.4), Inches(9.4), Inches(2.5)
    )
    tbl = tbl_shape.table

    # Row 0: months Jan-Jun
    for i, m in enumerate(months[:6]):
        tbl.cell(0, i + 1).text = m
    tbl.cell(0, 0).text = "Month"
    # Row 1: signups Jan-Jun
    tbl.cell(1, 0).text = "Signups"
    for i, s in enumerate(signups[:6]):
        tbl.cell(1, i + 1).text = str(s)
    # Row 2: months Jul-Dec header reuse
    tbl.cell(2, 0).text = "Month"
    for i, m in enumerate(months[6:]):
        tbl.cell(2, i + 1).text = f"{m}: {signups[6 + i]}"

    # ============================================================
    # Slide 3: "Growth Strategy" with bullet points
    # ============================================================
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Growth Strategy"
    p3.alignment = PP_ALIGN.LEFT
    run3 = p3.runs[0]
    run3.font.size = Pt(28)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    bullets = [
        "Expand referral program to increase organic acquisition by 25%",
        "Launch targeted social media campaigns in Q2 for emerging markets",
        "Optimize onboarding flow to improve Day-1 retention from 42% to 60%",
        "Partner with 3 industry influencers for co-branded content series",
        "Implement A/B testing framework for landing page conversion optimization",
    ]
    body3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5))
    btf = body3.text_frame
    btf.word_wrap = True
    for idx, b in enumerate(bullets):
        if idx == 0:
            par = btf.paragraphs[0]
        else:
            par = btf.add_paragraph()
        par.text = b
        par.level = 0
        for r in par.runs:
            r.font.size = Pt(18)
            r.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # ============================================================
    # Slide 4: "Cumulative User Growth" — title only, NO chart
    # ============================================================
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Cumulative User Growth"
    p4.alignment = PP_ALIGN.LEFT
    run4 = p4.runs[0]
    run4.font.size = Pt(28)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Subtitle note for context (the agent needs to create a chart here)
    sub4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(0.6))
    stf = sub4.text_frame
    sp = stf.paragraphs[0]
    sp.text = "Chart placeholder — visualize cumulative signups over 12 months"
    for r in sp.runs:
        r.font.size = Pt(14)
        r.font.italic = True
        r.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    # ============================================================
    # Slide 5: "Next Steps & Projections"
    # ============================================================
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf5 = txBox5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Next Steps & Projections"
    p5.alignment = PP_ALIGN.LEFT
    run5 = p5.runs[0]
    run5.font.size = Pt(28)
    run5.font.bold = True
    run5.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    next_items = [
        "Target 75,000 cumulative signups by end of Q1 2026",
        "Allocate $120K budget for paid acquisition channels",
        "Hire 2 additional growth engineers to support experimentation",
        "Deploy predictive churn model by March 2026",
    ]
    body5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(4))
    b5tf = body5.text_frame
    b5tf.word_wrap = True
    for idx, item in enumerate(next_items):
        if idx == 0:
            par = b5tf.paragraphs[0]
        else:
            par = b5tf.add_paragraph()
        par.text = item
        par.level = 0
        for r in par.runs:
            r.font.size = Pt(18)
            r.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
