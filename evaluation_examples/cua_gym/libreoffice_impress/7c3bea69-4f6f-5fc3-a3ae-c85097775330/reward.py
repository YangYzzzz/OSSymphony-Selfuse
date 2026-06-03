"""
Reward Script: Create area chart on slide 4 with semi-transparent blue fill
Task ID: impress_tct_053
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Area chart exists on slide 4
  Component 2 (0.25): Chart has 12 data points with upward trend
  Component 3 (0.25): Chart fill color is #1E88E5 (blue)
  Component 4 (0.20): Chart fill has ~40% opacity (semi-transparent)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_053'


def persist_app_state(domain):
    """Save any unsaved LibreOffice changes before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.chart import XL_CHART_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Verify we have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide4 = prs.slides[3]  # 0-indexed

    # Find chart shape on slide 4
    chart_shape = None
    for shape in slide4.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Area chart exists on slide 4 (0.30 points)
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            chart_type = chart.chart_type
            # AREA type has value 1 in XL_CHART_TYPE
            # Also accept AREA_STACKED (76) and AREA_STACKED_100 (77)
            area_types = {1, 76, 77}  # AREA, AREA_STACKED, AREA_STACKED_100
            actual_val = chart_type if isinstance(chart_type, int) else chart_type.value if hasattr(chart_type, 'value') else int(chart_type)
            if actual_val in area_types:
                print(f"PASS: Component 1 — Area chart found on slide 4, type={chart_type} (0.30 pts)")
                total_score += 0.30
            else:
                print(f"FAIL: Component 1 — Chart found but not area type, got type={chart_type}")
        else:
            print("FAIL: Component 1 — No chart found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # If no chart, remaining components cannot be checked
    if chart_shape is None:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    chart = chart_shape.chart

    # Component 2: Chart has 12 data points with upward trend (0.25 points)
    try:
        series_count = len(chart.series)
        if series_count >= 1:
            values = list(chart.series[0].values)
            num_points = len(values)
            categories = [str(c) for c in chart.plots[0].categories]
            num_cats = len(categories)

            points_ok = (num_points == 12)
            # Check upward trend: at least the last value > first value
            # and most consecutive pairs are increasing
            trend_ok = False
            if num_points >= 2:
                increasing_pairs = sum(1 for i in range(len(values)-1) if values[i+1] > values[i])
                # Allow some slack: at least 8 out of 11 pairs increasing
                trend_ok = (increasing_pairs >= 8) and (values[-1] > values[0])

            if points_ok and trend_ok:
                print(f"PASS: Component 2 — 12 data points with upward trend (first={values[0]}, last={values[-1]}, {increasing_pairs}/11 increasing) (0.25 pts)")
                total_score += 0.25
            elif points_ok:
                print(f"PARTIAL: Component 2 — 12 data points but weak trend ({increasing_pairs}/11 increasing)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 2 — Expected 12 data points, found {num_points}. Trend increasing pairs: {increasing_pairs if num_points >= 2 else 'N/A'}")
        else:
            print("FAIL: Component 2 — No series found in chart")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Chart fill color is #1E88E5 (blue) (0.25 points)
    try:
        series = chart.series[0]
        fill = series.format.fill
        fill_type = fill.type
        if fill_type is not None:
            try:
                rgb = str(fill.fore_color.rgb)
                if rgb.upper() == '1E88E5':
                    print(f"PASS: Component 3 — Fill color is #1E88E5 (0.25 pts)")
                    total_score += 0.25
                else:
                    print(f"FAIL: Component 3 — Fill color is #{rgb}, expected #1E88E5")
            except Exception as e:
                print(f"FAIL: Component 3 — Could not read fill color: {e}")
        else:
            print("FAIL: Component 3 — No fill type set on series")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Chart fill has ~40% opacity via XML alpha check (0.20 points)
    try:
        found_alpha = False
        target_alpha = 40000  # 40% in PowerPoint units (100000 = 100%)
        tolerance = 10000     # Allow 30%-50% range

        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        with zipfile.ZipFile(file_path, 'r') as zf:
            chart_files = [f for f in zf.namelist() if f.startswith('ppt/charts/') and f.endswith('.xml')]
            for cf in chart_files:
                content = zf.read(cf).decode('utf-8')
                root = ET.fromstring(content)
                # Find srgbClr elements with value 1E88E5 that have alpha child
                for elem in root.iter():
                    tag_local = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
                    if tag_local == 'srgbClr' and elem.get('val', '').upper() == '1E88E5':
                        for child in elem:
                            child_local = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                            if child_local == 'alpha':
                                alpha_val = int(child.get('val', '0'))
                                if abs(alpha_val - target_alpha) <= tolerance:
                                    found_alpha = True
                                    print(f"PASS: Component 4 — Alpha value {alpha_val} (~{alpha_val/1000:.0f}% opacity) within tolerance (0.20 pts)")
                                    total_score += 0.20
                                else:
                                    print(f"FAIL: Component 4 — Alpha value {alpha_val} ({alpha_val/1000:.0f}% opacity), expected ~40%")
                                break
                    if found_alpha:
                        break
                if found_alpha:
                    break

        if not found_alpha and total_score < 0.75:
            print("FAIL: Component 4 — No alpha/opacity setting found on chart fill with color #1E88E5")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
