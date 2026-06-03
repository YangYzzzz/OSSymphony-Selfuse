"""
Initial Setup: Create a 12-slide dissertation presentation
Task ID: impress_ndo_077
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_077'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide data: (title, bullet points)
    slides_data = [
        ("Introduction", [
            "Background on computational linguistics and NLP",
            "Research motivation: improving sentiment analysis accuracy",
            "Problem statement and research gap",
            "Thesis structure overview",
        ]),
        ("Literature Review", [
            "Evolution of sentiment analysis from lexicon-based to deep learning",
            "Transformer architectures: BERT, GPT, and domain-specific models",
            "Cross-domain sentiment transfer challenges (Pang & Lee, 2008)",
            "Gap: limited work on multi-modal sentiment in academic discourse",
        ]),
        ("Methodology", [
            "Mixed-methods research design combining quantitative and qualitative",
            "Phase 1: Corpus construction from 2,400 academic reviews",
            "Phase 2: Fine-tuning RoBERTa with domain-specific tokenization",
            "Phase 3: Human evaluation with 15 expert annotators",
        ]),
        ("Data Collection", [
            "Sources: PeerRead dataset, OpenReview API, ACL Anthology",
            "Time period: January 2019 to December 2024",
            "Total samples: 2,437 peer review documents",
            "Annotation scheme: 5-point Likert scale for sentiment polarity",
            "Inter-annotator agreement: Cohen's kappa = 0.82",
        ]),
        ("Analysis", [
            "Preprocessing: tokenization, stop-word removal, lemmatization",
            "Feature extraction using TF-IDF and contextual embeddings",
            "Model comparison: Logistic Regression, SVM, BiLSTM, RoBERTa",
            "Hyperparameter tuning via Bayesian optimization (100 trials)",
            "Cross-validation: stratified 5-fold on balanced dataset",
        ]),
        ("Results", [
            "RoBERTa achieved F1=0.891 on the test set (best overall)",
            "BiLSTM: F1=0.843, SVM: F1=0.796, LogReg: F1=0.771",
            "Ablation study: domain tokenizer improved F1 by +0.034",
            "Error analysis: sarcasm and hedging caused 62% of misclassifications",
        ]),
        ("Discussion", [
            "Domain-specific fine-tuning outperforms generic pre-training",
            "Contextual embeddings capture nuanced academic sentiment",
            "Comparison with prior work: +4.7% improvement over Wang et al. (2023)",
            "Practical implications for automated peer review assistance",
        ]),
        ("Implications", [
            "Theoretical: extends sentiment analysis frameworks to academic text",
            "Practical: prototype tool for conference program committees",
            "Methodological: reusable annotation guidelines for academic sentiment",
            "Policy: supports transparent and fair peer review processes",
        ]),
        ("Limitations", [
            "Dataset restricted to English-language computer science venues",
            "Limited generalizability to humanities and social sciences",
            "Model requires GPU resources not available in all settings",
            "Annotation subjectivity despite high inter-annotator agreement",
        ]),
        ("Future Work", [
            "Extend corpus to multilingual academic reviews",
            "Incorporate aspect-based sentiment for fine-grained feedback",
            "Explore few-shot learning to reduce annotation requirements",
            "Longitudinal study on model drift over time",
        ]),
        ("References", [
            "Devlin, J. et al. (2019). BERT: Pre-training of Deep Bidirectional Transformers",
            "Liu, Y. et al. (2019). RoBERTa: A Robustly Optimized BERT Pretraining",
            "Pang, B. & Lee, L. (2008). Opinion Mining and Sentiment Analysis",
            "Wang, X. et al. (2023). Sentiment in Scholarly Peer Review",
            "Vaswani, A. et al. (2017). Attention Is All You Need",
        ]),
        ("Appendix", [
            "A. Full annotation guidelines (12 pages)",
            "B. Hyperparameter search space and final configurations",
            "C. Additional confusion matrices for all model variants",
            "D. Survey instrument used for expert evaluation",
            "E. IRB approval documentation (Protocol #2024-0347)",
        ]),
    ]

    for title, bullets in slides_data:
        # Use Title+Content layout (index 1) for content slides
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title

        # Set title font
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)

        # Add bullet content to the body placeholder
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.space_after = Pt(6)
            for run in p.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
