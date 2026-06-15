"""
Reward Script: Add a Summary slide listing all slide titles as bullet points
Task ID: impress_ndo_077
Domain: libreoffice_impress
Scoring:
  Component 1: Presentation has 13 slides (0.2 pts)
  Component 2: Last slide title is 'Summary' (0.3 pts)
  Component 3: Summary slide body contains all 12 original titles (0.5 pts, progressive)
"""

import os

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_077'

EXPECTED_TITLES = [
    'Introduction',
    'Literature Review',
    'Methodology',
    'Data Collection',
    'Analysis',
    'Results',
    'Discussion',
    'Implications',
    'Limitations',
    'Future Work',
    'References',
    'Appendix',
]


def get_all_text_shapes(slide):
    """Recursively get all text shapes including those inside groups."""
    def extract(shape):
        results = []
        if hasattr(shape, "text") and hasattr(shape, "text_frame"):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.extend(extract(sub))
        return results
    out = []
    for shape in slide.shapes:
        out.extend(extract(shape))
    return out


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Presentation has 13 slides (0.2 points)
    # Initial has 12, golden has 13 — this checks the new slide was added
    try:
        if num_slides == 13:
            print(f"PASS: Component 1 -- Slide count is 13 (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 1 -- Expected 13 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Last slide title is 'Summary' (0.3 points)
    # Initial has no Summary slide; golden has it as slide 13
    try:
        if num_slides < 13:
            print(f"FAIL: Component 2 -- Not enough slides to check last slide title")
        else:
            last_slide = prs.slides[num_slides - 1]
            last_title = last_slide.shapes.title.text.strip() if last_slide.shapes.title else ''
            if last_title.lower() == 'summary':
                print(f"PASS: Component 2 -- Last slide title is '{last_title}' (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 2 -- Expected last slide title 'Summary', found '{last_title}'")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Summary slide body contains all 12 original titles as bullet points (0.5 points)
    # Progressive: each matched title earns 0.5/12 points
    # Initial has no summary slide so this scores 0.0 on initial
    try:
        if num_slides < 13:
            print(f"FAIL: Component 3 -- No summary slide to check bullet points")
        else:
            last_slide = prs.slides[num_slides - 1]
            # Collect all non-title text from the summary slide
            text_shapes = get_all_text_shapes(last_slide)
            body_texts = []
            title_text = ''
            if last_slide.shapes.title:
                title_text = last_slide.shapes.title.text.strip()

            for shape in text_shapes:
                for para in shape.text_frame.paragraphs:
                    txt = para.text.strip()
                    # Skip the title text and empty lines
                    if txt and txt != title_text:
                        body_texts.append(txt)

            print(f"  Summary slide body texts found: {body_texts}")

            matched = 0
            per_title_score = 0.5 / 12.0
            for expected in EXPECTED_TITLES:
                # Check if any body text line contains this title (case-insensitive)
                found = any(expected.lower() in bt.lower() for bt in body_texts)
                if found:
                    matched += 1
                    print(f"  MATCH: '{expected}' found in summary bullet points")
                else:
                    print(f"  MISS: '{expected}' NOT found in summary bullet points")

            comp3_score = matched * per_title_score
            if matched == 12:
                print(f"PASS: Component 3 -- All 12 titles found as bullet points ({comp3_score:.3f} pts)")
                total_score += comp3_score
            elif matched > 0:
                print(f"PARTIAL: Component 3 -- {matched}/12 titles found ({comp3_score:.3f} pts)")
                total_score += comp3_score
            else:
                print(f"FAIL: Component 3 -- 0/12 titles found in summary slide")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.3f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
