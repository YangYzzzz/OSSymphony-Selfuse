"""
Reward Script: Supply-Demand Diagram on Slide 7
Task ID: impress_stu_086
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): X and Y axes present with labels (Quantity, Price)
  Component 2 (0.20): Demand curve (red line) and Supply curve (blue line) present
  Component 3 (0.15): Equilibrium point shape present with green fill
  Component 4 (0.15): Equilibrium label and P*/Q* text box present
  Component 5 (0.30): Animation sequence exists on slide 7
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_086'


def persist_app_state(domain):
    """Try to save any open LibreOffice document before verifying."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print(f"PERSIST: ctrl+s sent for {domain}")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.oxml.ns import qn
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 7 slides
    if len(prs.slides) < 7:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 7")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[6]  # Slide 7 (0-indexed)
    shapes = list(slide.shapes)

    # Helper: find shapes by type and text content
    def get_shape_names_and_types():
        result = []
        for s in shapes:
            text = ""
            if s.has_text_frame:
                text = " ".join(p.text for p in s.text_frame.paragraphs).strip()
            result.append({
                "name": s.name,
                "type": s.shape_type,
                "text": text,
                "shape_id": s.shape_id,
                "element": s.element,
                "shape": s,
            })
        return result

    shape_info = get_shape_names_and_types()

    # Helper: get line color from shape element
    def get_line_color(el):
        try:
            for ln in el.findall('.//' + qn('a:ln')):
                sf = ln.find(qn('a:solidFill'))
                if sf is not None:
                    srgb = sf.find(qn('a:srgbClr'))
                    if srgb is not None:
                        return srgb.attrib.get('val', '').upper()
        except:
            pass
        return None

    # Helper: get fill color from shape element
    def get_fill_color(el):
        try:
            for spPr_tag in [qn('p:spPr'), qn('a:spPr')]:
                spPr = el.find(spPr_tag)
                if spPr is not None:
                    sf = spPr.find(qn('a:solidFill'))
                    if sf is not None:
                        srgb = sf.find(qn('a:srgbClr'))
                        if srgb is not None:
                            return srgb.attrib.get('val', '').upper()
            # Also search deeper
            for sf in el.findall('.//' + qn('a:solidFill')):
                srgb = sf.find(qn('a:srgbClr'))
                if srgb is not None:
                    return srgb.attrib.get('val', '').upper()
        except:
            pass
        return None

    # Find line shapes and text shapes on slide 7
    line_shapes = [s for s in shape_info if s["type"] == MSO_SHAPE_TYPE.LINE]
    text_shapes = [s for s in shape_info if s["type"] == MSO_SHAPE_TYPE.TEXT_BOX]
    auto_shapes = [s for s in shape_info if s["type"] == MSO_SHAPE_TYPE.AUTO_SHAPE]

    # ---------------------------------------------------------------
    # Component 1: X and Y axes with labels (0.20 points)
    # Task requires axes with labels 'Quantity' (X) and 'Price' (Y)
    # Initial state has NO axes and NO such labels, so this scores a change.
    # ---------------------------------------------------------------
    try:
        # Check for at least 2 line shapes (axes) beyond the initial 2 shapes
        axes_found = 0
        for ls in line_shapes:
            name_lower = ls["name"].lower()
            if "axis" in name_lower or "axe" in name_lower:
                axes_found += 1
        # Also accept generic lines that could be axes (horizontal/vertical)
        if axes_found < 2:
            # Count lines that are perfectly horizontal or vertical
            for ls in line_shapes:
                s = ls["shape"]
                if s.width <= 1 or s.height <= 1:
                    axes_found += 1

        # Check for Quantity and Price labels
        quantity_label = any("quantity" in s["text"].lower() for s in text_shapes)
        price_label = any("price" in s["text"].lower() for s in text_shapes)

        if axes_found >= 2 and quantity_label and price_label:
            print(f"PASS: Component 1 -- Axes ({axes_found}) with Quantity and Price labels (0.20 pts)")
            total_score += 0.20
        elif axes_found >= 2 and (quantity_label or price_label):
            print(f"PARTIAL: Component 1 -- Axes found but missing a label (0.10 pts)")
            total_score += 0.10
        elif axes_found >= 1:
            print(f"PARTIAL: Component 1 -- Only {axes_found} axis found (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 1 -- No axes found. Lines: {len(line_shapes)}, Quantity label: {quantity_label}, Price label: {price_label}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # ---------------------------------------------------------------
    # Component 2: Demand curve (red) and Supply curve (blue) (0.20 points)
    # Task requires red demand curve and blue supply curve.
    # Initial state has NO such colored lines.
    # ---------------------------------------------------------------
    try:
        red_line_found = False
        blue_line_found = False
        for ls in line_shapes:
            color = get_line_color(ls["element"])
            name_lower = ls["name"].lower()
            if color == "FF0000" or "demand" in name_lower:
                if color == "FF0000":
                    red_line_found = True
            if color == "0000FF" or "supply" in name_lower:
                if color == "0000FF":
                    blue_line_found = True

        # Also check for named demand/supply even if colors differ slightly
        demand_by_name = any("demand" in ls["name"].lower() for ls in line_shapes)
        supply_by_name = any("supply" in ls["name"].lower() for ls in line_shapes)

        if red_line_found and blue_line_found:
            print(f"PASS: Component 2 -- Demand (red) and Supply (blue) curves found (0.20 pts)")
            total_score += 0.20
        elif (red_line_found or demand_by_name) and (blue_line_found or supply_by_name):
            print(f"PARTIAL: Component 2 -- Curves found but color mismatch (0.12 pts)")
            total_score += 0.12
        elif red_line_found or blue_line_found or demand_by_name or supply_by_name:
            print(f"PARTIAL: Component 2 -- Only one curve found (0.08 pts)")
            total_score += 0.08
        else:
            print(f"FAIL: Component 2 -- No demand/supply curves found")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # ---------------------------------------------------------------
    # Component 3: Equilibrium point shape (0.15 points)
    # Task requires a point at intersection labeled 'Equilibrium'.
    # Initial state has NO such shape (only a text box that python-pptx
    # may report as AUTO_SHAPE containing 'Market Equilibrium').
    # We filter out any auto_shape whose text contains 'market equilibrium'
    # to avoid false positives on the pre-existing title text box.
    # ---------------------------------------------------------------
    try:
        eq_point_found = False
        eq_point_colored = False

        # Filter out auto shapes that are really the title/header text box
        candidate_auto_shapes = [
            s for s in auto_shapes
            if "market equilibrium" not in s["text"].lower()
        ]

        for ashape in candidate_auto_shapes:
            name_lower = ashape["name"].lower()
            if "equilibrium" in name_lower or "point" in name_lower:
                eq_point_found = True
                fill = get_fill_color(ashape["element"])
                if fill is not None:
                    eq_point_colored = True

        # Fallback: any new auto shape (not the title text) with a fill
        if not eq_point_found and len(candidate_auto_shapes) > 0:
            eq_point_found = True
            fill = get_fill_color(candidate_auto_shapes[0]["element"])
            if fill is not None:
                eq_point_colored = True

        if eq_point_found and eq_point_colored:
            print(f"PASS: Component 3 -- Equilibrium point with fill color found (0.15 pts)")
            total_score += 0.15
        elif eq_point_found:
            print(f"PARTIAL: Component 3 -- Equilibrium point found but no fill (0.08 pts)")
            total_score += 0.08
        else:
            print(f"FAIL: Component 3 -- No equilibrium point shape found (candidates: {len(candidate_auto_shapes)})")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # ---------------------------------------------------------------
    # Component 4: Equilibrium label + P*/Q* text (0.15 points)
    # Task requires 'Equilibrium' label and 'P* = $15, Q* = 100' text box.
    # Initial state has NO such text boxes.
    # ---------------------------------------------------------------
    try:
        eq_label_found = False
        pq_text_found = False

        for ts in text_shapes:
            text_lower = ts["text"].lower()
            # Check for equilibrium label (not the title "Market Equilibrium")
            if "equilibrium" in text_lower and "market" not in text_lower:
                eq_label_found = True
            # Check for P* and Q* values
            if "p*" in text_lower and "q*" in text_lower:
                pq_text_found = True
            elif "$15" in ts["text"] and "100" in ts["text"]:
                pq_text_found = True

        if eq_label_found and pq_text_found:
            print(f"PASS: Component 4 -- Equilibrium label and P*/Q* text found (0.15 pts)")
            total_score += 0.15
        elif eq_label_found or pq_text_found:
            found = "Equilibrium label" if eq_label_found else "P*/Q* text"
            print(f"PARTIAL: Component 4 -- Only {found} found (0.08 pts)")
            total_score += 0.08
        else:
            print(f"FAIL: Component 4 -- No equilibrium label or P*/Q* text found")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # ---------------------------------------------------------------
    # Component 5: Animation sequence on slide 7 (0.30 points)
    # Task requires animations: axes appear, demand wipes left, supply wipes right,
    # equilibrium grows, P*/Q* fades in.
    # Initial state has NO animations on slide 7.
    # ---------------------------------------------------------------
    try:
        ns = {
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        }

        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide7.xml') as f:
                root = ET.parse(f).getroot()
                timing = root.find('.//p:timing', ns)

                if timing is not None:
                    # Count animation entries (par elements with presetClass="entr")
                    anim_entries = timing.findall('.//' + '{http://schemas.openxmlformats.org/presentationml/2006/main}par')
                    entry_anims = [a for a in anim_entries if a.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cTn') is not None]

                    # Count cTn with presetClass="entr"
                    entrance_anims = timing.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cTn[@presetClass="entr"]')
                    num_entrance = len(entrance_anims)

                    # Check for wipe animations (presetID=22)
                    wipe_anims = [a for a in entrance_anims if a.get('presetID') == '22']
                    # Check for grow/scale animations (presetID=53)
                    grow_anims = [a for a in entrance_anims if a.get('presetID') == '53']
                    # Check for appear animations (presetID=1)
                    appear_anims = [a for a in entrance_anims if a.get('presetID') == '1']

                    anim_score = 0.0

                    # Sub-check 5a: Any animations exist at all (0.10)
                    if num_entrance >= 3:
                        print(f"  PASS: Component 5a -- {num_entrance} entrance animations found (0.10 pts)")
                        anim_score += 0.10
                    elif num_entrance >= 1:
                        print(f"  PARTIAL: Component 5a -- Only {num_entrance} entrance animations (0.05 pts)")
                        anim_score += 0.05
                    else:
                        print(f"  FAIL: Component 5a -- No entrance animations found")

                    # Sub-check 5b: Wipe animations for demand/supply curves (0.10)
                    if len(wipe_anims) >= 2:
                        print(f"  PASS: Component 5b -- {len(wipe_anims)} wipe animations found (0.10 pts)")
                        anim_score += 0.10
                    elif len(wipe_anims) >= 1:
                        print(f"  PARTIAL: Component 5b -- Only {len(wipe_anims)} wipe animation (0.05 pts)")
                        anim_score += 0.05
                    else:
                        print(f"  FAIL: Component 5b -- No wipe animations found")

                    # Sub-check 5c: Grow animation for equilibrium + appear/fade for others (0.10)
                    if len(grow_anims) >= 1 and len(appear_anims) >= 2:
                        print(f"  PASS: Component 5c -- Grow ({len(grow_anims)}) and appear ({len(appear_anims)}) animations found (0.10 pts)")
                        anim_score += 0.10
                    elif len(grow_anims) >= 1 or len(appear_anims) >= 2:
                        print(f"  PARTIAL: Component 5c -- Grow: {len(grow_anims)}, Appear: {len(appear_anims)} (0.05 pts)")
                        anim_score += 0.05
                    else:
                        print(f"  FAIL: Component 5c -- No grow or appear animations found")

                    total_score += anim_score
                    print(f"PASS: Component 5 -- Animation sequence ({anim_score:.2f} pts)")
                else:
                    print(f"FAIL: Component 5 -- No timing/animation data found on slide 7")

    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
