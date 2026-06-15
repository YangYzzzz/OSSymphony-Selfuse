"""
Initial Setup: Create a 10-slide Microecon 101 presentation with slide 7 empty (titled 'Market Equilibrium').
Task ID: impress_stu_086
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_086'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def add_blank_titled_slide(prs, title_text):
    """Add a slide with only a title (Title Only layout)."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    add_title_slide(prs, "Microeconomics 101", "Principles of Supply, Demand, and Market Structure\nDr. Rebecca Thornton | Fall 2025")

    # Slide 2: Course Overview
    add_content_slide(prs, "Course Overview", [
        "Introduction to fundamental microeconomic concepts",
        "Understanding consumer and producer behavior",
        "Market structures: perfect competition to monopoly",
        "Policy implications and welfare analysis",
        "Real-world case studies from technology and healthcare sectors",
    ])

    # Slide 3: What is Economics?
    add_content_slide(prs, "What is Economics?", [
        "The study of how societies allocate scarce resources",
        "Microeconomics focuses on individual agents and markets",
        "Key assumption: rational agents maximize utility or profit",
        "Positive analysis (what is) vs. normative analysis (what should be)",
        "Ceteris paribus: holding other factors constant",
    ])

    # Slide 4: Demand
    add_content_slide(prs, "The Law of Demand", [
        "As price increases, quantity demanded decreases (ceteris paribus)",
        "Demand curve slopes downward from left to right",
        "Determinants: income, preferences, prices of substitutes/complements",
        "Individual demand vs. market demand (horizontal summation)",
        "Example: When coffee price rises from $3 to $5, daily cups sold drop from 200 to 120",
    ])

    # Slide 5: Supply
    add_content_slide(prs, "The Law of Supply", [
        "As price increases, quantity supplied increases (ceteris paribus)",
        "Supply curve slopes upward from left to right",
        "Determinants: input costs, technology, number of producers",
        "Short-run vs. long-run supply elasticity",
        "Example: A 10% rise in wheat prices leads to 6% increase in quantity supplied",
    ])

    # Slide 6: Elasticity
    add_content_slide(prs, "Price Elasticity", [
        "Elasticity = % change in quantity / % change in price",
        "Elastic demand (|E| > 1): luxury goods, many substitutes",
        "Inelastic demand (|E| < 1): necessities, few substitutes",
        "Unit elastic (|E| = 1): total revenue unchanged with price change",
        "Cross-price and income elasticity extend the concept further",
    ])

    # Slide 7: Market Equilibrium — EMPTY (only title)
    add_blank_titled_slide(prs, "Market Equilibrium")

    # Slide 8: Consumer Surplus
    add_content_slide(prs, "Consumer and Producer Surplus", [
        "Consumer surplus: difference between willingness to pay and actual price",
        "Producer surplus: difference between actual price and minimum acceptable price",
        "Total surplus = consumer surplus + producer surplus",
        "Equilibrium maximizes total surplus (allocative efficiency)",
        "Deadweight loss occurs when markets deviate from equilibrium",
    ])

    # Slide 9: Market Failures
    add_content_slide(prs, "Market Failures", [
        "Externalities: costs or benefits imposed on third parties",
        "Public goods: non-excludable and non-rivalrous consumption",
        "Information asymmetry: adverse selection and moral hazard",
        "Market power: monopolies and oligopolies restrict output",
        "Government intervention: taxes, subsidies, regulation, antitrust",
    ])

    # Slide 10: Summary
    add_content_slide(prs, "Key Takeaways", [
        "Supply and demand determine market prices and quantities",
        "Equilibrium occurs where supply meets demand (P* = $15, Q* = 100 in our example)",
        "Elasticity measures responsiveness to price changes",
        "Surplus analysis helps evaluate market efficiency",
        "Next lecture: Market structures and competitive strategy",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
