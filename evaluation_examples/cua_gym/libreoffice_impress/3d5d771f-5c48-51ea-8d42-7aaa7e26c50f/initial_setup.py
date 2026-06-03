"""
Initial Setup: team_presentation.pptx with blank slide 2 (title only) and team_photo.jpg
Task ID: impress_media_032
Domain: libreoffice_impress

Creates:
  - ~/Desktop/team_presentation.pptx  (3 slides; slide 2 has title but NO image)
  - ~/Pictures/team_photo.jpg          (1200x800 px realistic team photo placeholder)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont
import io

DESKTOP = '/home/user/Desktop'
PICTURES = '/home/user/Pictures'
TASK_ID = 'impress_media_032'
PPTX_PATH = f'{DESKTOP}/team_presentation.pptx'
IMG_PATH = f'{PICTURES}/team_photo.jpg'


# ─────────────────────────────────────────────
# 1.  Create team_photo.jpg  (1200×800 px)
# ─────────────────────────────────────────────
def create_team_photo():
    os.makedirs(PICTURES, exist_ok=True)

    img = Image.new('RGB', (1200, 800), color=(70, 110, 160))
    draw = ImageDraw.Draw(img)

    # Sky gradient background (simple top-to-bottom)
    for y in range(800):
        r = int(100 + (y / 800) * 60)
        g = int(140 + (y / 800) * 40)
        b = int(200 - (y / 800) * 50)
        draw.line([(0, y), (1200, y)], fill=(r, g, b))

    # Ground area
    draw.rectangle([0, 560, 1200, 800], fill=(80, 120, 60))

    # Draw simplified silhouettes of 6 team members
    team_positions = [100, 250, 400, 550, 750, 900, 1050]
    body_color = (40, 40, 80)
    skin_tones = [
        (210, 160, 120), (180, 120, 80), (240, 200, 170),
        (160, 100, 60),  (220, 175, 140), (190, 140, 100),
        (230, 190, 155),
    ]
    shirt_colors = [
        (200, 50, 50), (50, 100, 200), (50, 180, 80),
        (220, 160, 30), (120, 50, 180), (50, 180, 180),
        (200, 100, 50),
    ]

    for i, x in enumerate(team_positions):
        skin = skin_tones[i % len(skin_tones)]
        shirt = shirt_colors[i % len(shirt_colors)]
        # Head
        draw.ellipse([x - 20, 430, x + 20, 470], fill=skin)
        # Body (shirt)
        draw.rectangle([x - 25, 470, x + 25, 560], fill=shirt)
        # Arms
        draw.rectangle([x - 40, 480, x - 25, 540], fill=shirt)
        draw.rectangle([x + 25, 480, x + 40, 540], fill=shirt)
        # Legs
        draw.rectangle([x - 20, 560, x - 5,  620], fill=(30, 30, 80))
        draw.rectangle([x + 5,  560, x + 20, 620], fill=(30, 30, 80))
        # Smile
        draw.arc([x - 10, 447, x + 10, 462], start=0, end=180, fill=(80, 40, 20), width=2)

    # Company banner at top
    draw.rectangle([200, 30, 1000, 90], fill=(255, 255, 255, 200))
    draw.rectangle([200, 30, 1000, 90], outline=(30, 60, 120), width=3)

    # Text labels (using default font since no TTF available easily)
    try:
        font_large = ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf', 32)
        font_small = ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf', 22)
    except Exception:
        font_large = ImageFont.load_default()
        font_small = font_large

    draw.text((600, 60), 'Nexus Solutions — Annual Team Day 2025',
              fill=(30, 60, 120), font=font_large, anchor='mm')
    draw.text((600, 720), 'Engineering · Marketing · Operations · Design',
              fill=(240, 240, 240), font=font_small, anchor='mm')

    img.save(IMG_PATH, 'JPEG', quality=92)
    print(f'Team photo created: {IMG_PATH}  ({img.size[0]}×{img.size[1]} px)')


# ─────────────────────────────────────────────
# 2.  Create team_presentation.pptx
# ─────────────────────────────────────────────
def create_presentation():
    os.makedirs(DESKTOP, exist_ok=True)

    prs = Presentation()
    # Standard widescreen 16:9
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title slide ---
    sl1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
    sl1.shapes.title.text = 'Nexus Solutions'
    sl1.placeholders[1].text = 'Annual Company Review 2025'
    # Style title
    title_run = sl1.shapes.title.text_frame.paragraphs[0].runs[0]
    title_run.font.size = Pt(44)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0x1F, 0x3C, 0x78)

    subtitle_run = sl1.placeholders[1].text_frame.paragraphs[0].runs[0]
    subtitle_run.font.size = Pt(28)
    subtitle_run.font.color.rgb = RGBColor(0x44, 0x72, 0xC4)

    # Background color for slide 1
    fill1 = sl1.background.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFB)

    # --- Slide 2: Meet the Team  (title only, NO image) ---
    # Use layout 5 (Blank) to avoid placeholder issues, add title text box manually
    sl2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add a title text box at the top
    title_box = sl2.shapes.add_textbox(
        Inches(1), Inches(0.3), Inches(11.33), Inches(1.2)
    )
    title2_tf = title_box.text_frame
    title2_tf.word_wrap = True
    para2 = title2_tf.paragraphs[0]
    para2.text = 'Meet the Team'
    para2.alignment = PP_ALIGN.CENTER
    run2 = para2.runs[0]
    run2.font.size = Pt(40)
    run2.font.bold = True
    run2.font.color.rgb = RGBColor(0x1F, 0x3C, 0x78)

    # Subtle background
    fill2 = sl2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFF)

    # --- Slide 3: Our Achievements ---
    sl3 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
    if sl3.shapes.title is not None:
        sl3.shapes.title.text = 'Our 2025 Achievements'
        if sl3.shapes.title.text_frame.paragraphs[0].runs:
            title3_run = sl3.shapes.title.text_frame.paragraphs[0].runs[0]
            title3_run.font.size = Pt(36)
            title3_run.font.bold = True
            title3_run.font.color.rgb = RGBColor(0x1F, 0x3C, 0x78)
    tf3 = sl3.placeholders[1].text_frame
    tf3.text = 'Revenue grew 32% year-over-year'
    bullets = [
        'Launched 4 new product lines across APAC markets',
        'Expanded headcount from 85 to 124 employees',
        'Achieved ISO 27001 security certification',
        'Opened new offices in Singapore and Sydney',
        'Customer satisfaction score reached 4.8 / 5.0',
    ]
    for bullet in bullets:
        p = tf3.add_paragraph()
        p.text = bullet
        p.level = 1
        if p.runs:
            p.runs[0].font.size = Pt(20)

    fill3 = sl3.background.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFB)

    # --- Slide 4: Looking Ahead ---
    sl4 = prs.slides.add_slide(prs.slide_layouts[1])
    if sl4.shapes.title is not None:
        sl4.shapes.title.text = 'Looking Ahead: 2026 Goals'
        if sl4.shapes.title.text_frame.paragraphs[0].runs:
            title4_run = sl4.shapes.title.text_frame.paragraphs[0].runs[0]
            title4_run.font.size = Pt(36)
            title4_run.font.bold = True
            title4_run.font.color.rgb = RGBColor(0x1F, 0x3C, 0x78)
    tf4 = sl4.placeholders[1].text_frame
    tf4.text = 'Strategic priorities for the coming year'
    goals = [
        'Enter three new international markets',
        'Release AI-powered analytics dashboard',
        'Grow partnerships with Fortune 500 companies',
        'Achieve carbon-neutral operations by Q3',
        'Scale engineering team to 200 professionals',
    ]
    for goal in goals:
        p = tf4.add_paragraph()
        p.text = goal
        p.level = 1
        if p.runs:
            p.runs[0].font.size = Pt(20)

    fill4 = sl4.background.fill
    fill4.solid()
    fill4.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFB)

    prs.save(PPTX_PATH)
    print(f'Presentation created: {PPTX_PATH}')
    print(f'  Slides: {len(prs.slides)}')
    print(f'  Slide 2: "Meet the Team" title textbox, NO image (blank body)')

    # Also save a copy as the canonical initial file for adversarial pipeline
    import shutil
    initial_copy = f'/home/user/{TASK_ID}_initial.pptx'
    shutil.copy(PPTX_PATH, initial_copy)
    print(f'Initial copy saved: {initial_copy}')


create_team_photo()
create_presentation()
print('initial_setup.py complete.')
