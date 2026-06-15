"""
Reward Script: Insert team_photo.jpg into slide 2, crop to center 60% horizontally,
               center on slide, and add a drop shadow effect.
Task ID: impress_media_032
Domain: libreoffice_impress (python-pptx)
Scoring:
  Component 1: Image inserted on slide 2            — 0.25 pts
  Component 2: Horizontal crop 20% left + 20% right — 0.30 pts
  Component 3: Image horizontally centered on slide  — 0.25 pts
  Component 4: Drop shadow effect applied            — 0.20 pts
  Total: 1.0
"""

import os
import zipfile
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'  # VM path — reward scripts run on the VM
TASK_ID = 'impress_media_032'

CROP_TOLERANCE = 0.01   # ±1% tolerance for crop fraction checks
POS_TOLERANCE  = 0.005  # 0.5% relative tolerance for position checks


def is_approx_equal(val1, val2, tolerance=POS_TOLERANCE):
    """Return True if val1 and val2 are within relative tolerance of each other."""
    if val1 == val2:
        return True
    denom = max(abs(val1), abs(val2))
    if denom == 0:
        return True
    return abs(val1 - val2) / denom <= tolerance


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # --- Load presentation ---
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"CRITICAL: Expected at least 2 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]  # 0-indexed; slide 2 is index 1
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Find picture shape(s) on slide 2
    picture_shapes = [s for s in slide2.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]

    # -----------------------------------------------------------------------
    # Component 1: Image is inserted on slide 2 (0.25 points)
    # FAILS on initial (no picture on slide 2), PASSES on golden (1 picture)
    # -----------------------------------------------------------------------
    try:
        if len(picture_shapes) >= 1:
            pic = picture_shapes[0]
            print(f"PASS: Component 1 — Image found on slide 2 (shape name: {pic.name}) (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — No picture found on slide 2 (found {len(picture_shapes)} pictures)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # For components 2–4, we need the picture to exist
    if len(picture_shapes) == 0:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    pic = picture_shapes[0]

    # -----------------------------------------------------------------------
    # Component 2: Horizontal crop is 20% from left AND 20% from right (0.30 pts)
    # Showing center 60% — crop_left ≈ 0.20, crop_right ≈ 0.20
    # crop_top and crop_bottom should be 0.0 (no vertical crop)
    # FAILS on initial (no crop at all), PASSES on golden.
    # -----------------------------------------------------------------------
    try:
        crop_left  = pic.crop_left
        crop_right = pic.crop_right
        crop_top   = pic.crop_top
        crop_bottom = pic.crop_bottom

        left_ok   = abs(crop_left  - 0.20) <= CROP_TOLERANCE
        right_ok  = abs(crop_right - 0.20) <= CROP_TOLERANCE
        top_ok    = abs(crop_top)           <= CROP_TOLERANCE
        bottom_ok = abs(crop_bottom)        <= CROP_TOLERANCE

        if left_ok and right_ok and top_ok and bottom_ok:
            print(f"PASS: Component 2 — Crop matches: left={crop_left:.4f}, right={crop_right:.4f}, "
                  f"top={crop_top:.4f}, bottom={crop_bottom:.4f} (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 2 — Crop mismatch: left={crop_left:.4f} (want 0.20±{CROP_TOLERANCE}), "
                  f"right={crop_right:.4f} (want 0.20±{CROP_TOLERANCE}), "
                  f"top={crop_top:.4f} (want 0.00), bottom={crop_bottom:.4f} (want 0.00)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # -----------------------------------------------------------------------
    # Component 3: Image is horizontally centered on the slide (0.25 points)
    # Centered means: pic.left == (slide_width - pic.width) / 2  (within 0.5%)
    # FAILS on initial (no picture), PASSES on golden.
    # -----------------------------------------------------------------------
    try:
        expected_left = (slide_width - pic.width) / 2
        actual_left   = pic.left

        if is_approx_equal(actual_left, expected_left):
            print(f"PASS: Component 3 — Image horizontally centered: left={actual_left} "
                  f"(expected≈{expected_left:.0f}) (0.25 pts)")
            total_score += 0.25
        else:
            offset_px = actual_left - expected_left
            print(f"FAIL: Component 3 — Image not centered: left={actual_left}, "
                  f"expected≈{expected_left:.0f}, offset={offset_px:.0f} EMU")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # -----------------------------------------------------------------------
    # Component 4: Drop shadow effect applied to the image (0.20 points)
    # Verified via ZIP/XML: <a:effectLst><a:outerShdw ...> present in p:spPr
    # FAILS on initial (no picture so no shadow), PASSES on golden.
    # -----------------------------------------------------------------------
    try:
        shadow_found = False
        pic_ns = 'http://schemas.openxmlformats.org/drawingml/2006/picture'
        a_ns   = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        with zipfile.ZipFile(file_path, 'r') as zf:
            # Slide 2 is slide index 1 → slide2.xml
            with zf.open('ppt/slides/slide2.xml') as f:
                root = ET.fromstring(f.read())

            # The picture element is under p: (presentationml) namespace in this file.
            # Use a broad search: find any outerShdw element anywhere in the slide XML.
            # This is valid because outerShdw only appears on shapes that have shadows.
            outer_shadows = root.findall(f'.//{{{a_ns}}}outerShdw')
            shadow_found = len(outer_shadows) > 0

        if shadow_found:
            print(f"PASS: Component 4 — Drop shadow (outerShdw) found on picture shape (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 4 — No outerShdw element found in effectLst on slide 2 picture")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # -----------------------------------------------------------------------
    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against golden file (path on VM)
file_path = f'{WORKDIR}/{TASK_ID}_initial.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
