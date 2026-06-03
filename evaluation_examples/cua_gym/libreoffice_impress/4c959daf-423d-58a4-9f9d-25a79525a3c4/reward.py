"""
Reward Script: Branded slide master with navy background, gold line, footer, and slide numbers
Task ID: impress_rp_002
Domain: libreoffice_impress
Scoring:
  Component 1: Master background is navy #1B2A4A (0.25 pts)
  Component 2: Gold horizontal line at y~7.0in, full width, stroke #C9A84C (0.30 pts)
  Component 3: Footer placeholder with 'Apex Dynamics' in white 10pt, bottom-left (0.25 pts)
  Component 4: Slide number placeholder positioned bottom-right (0.20 pts)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_002'


def is_approximately_equal(val1, val2, tolerance=0.02):
    """Check if two values are approximately equal within tolerance (relative)."""
    if val1 == val2:
        return True
    if val1 == 0 or val2 == 0:
        return abs(val1 - val2) < 50000  # allow ~50k EMU slack for zero comparisons
    return abs(val1 - val2) / max(abs(val1), abs(val2)) <= tolerance


def persist_app_state(domain):
    """Save any unsaved changes in LibreOffice."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(1.0)
            print("PERSIST: ctrl+s sent for libreoffice_impress")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Get slide master
    try:
        master = prs.slide_masters[0]
    except Exception as e:
        print(f"CRITICAL: Cannot access slide master: {e}")
        print("REWARD: 0.0")
        return 0.0

    slide_width = prs.slide_width  # expected ~12191695 EMU

    # Component 1: Master background is navy #1B2A4A (0.25 points)
    try:
        bg_fill = master.background.fill
        if bg_fill.type == 1:  # SOLID fill
            bg_color = str(bg_fill.fore_color.rgb).upper()
            if bg_color == "1B2A4A":
                print(f"PASS: Component 1 — Master background is #{bg_color} (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 — Master background is #{bg_color}, expected #1B2A4A")
        else:
            print(f"FAIL: Component 1 — Master background fill type is {bg_fill.type}, expected SOLID (1)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Gold horizontal line shape on master (0.30 points)
    # Looking for a shape at y~7.0in (6400800 EMU), spanning full width, with stroke #C9A84C
    try:
        gold_line_found = False
        gold_line_correct_stroke = False
        gold_line_correct_position = False

        for shape in master.shapes:
            # Skip placeholders
            if shape.is_placeholder:
                continue

            # Look for a thin shape near y=7.0in (6400800 EMU)
            # It should span full slide width and be thin (line-like)
            shape_top = shape.top
            shape_height = shape.height
            shape_width = shape.width
            shape_left = shape.left

            # Check if it's at approximately y=7.0in and spans full width
            near_y_7in = is_approximately_equal(shape_top, 6400800, tolerance=0.05)
            is_thin = shape_height < 100000  # less than ~0.1 inches
            spans_width = is_approximately_equal(shape_width, slide_width, tolerance=0.05)
            starts_left = shape_left < 200000  # starts near left edge

            if near_y_7in and is_thin and spans_width and starts_left:
                gold_line_found = True
                # Check stroke color via XML
                el = shape._element
                ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                ln = el.find('.//a:ln', ns)
                if ln is not None:
                    sf = ln.find('.//a:solidFill', ns)
                    if sf is not None:
                        srgb = sf.find('a:srgbClr', ns)
                        if srgb is not None:
                            color_val = srgb.get('val', '').upper()
                            if color_val == 'C9A84C':
                                gold_line_correct_stroke = True
                                print(f"  Line stroke color: #{color_val} -- correct")
                            else:
                                print(f"  Line stroke color: #{color_val} -- expected #C9A84C")

                gold_line_correct_position = True
                print(f"  Line position: top={shape_top}, left={shape_left}, w={shape_width}, h={shape_height}")
                break

        if gold_line_found and gold_line_correct_stroke and gold_line_correct_position:
            print(f"PASS: Component 2 — Gold line found with correct position and stroke (0.30 pts)")
            total_score += 0.30
        elif gold_line_found and gold_line_correct_position:
            print(f"PARTIAL: Component 2 — Line found at correct position but wrong stroke color (0.15 pts)")
            total_score += 0.15
        elif gold_line_found:
            print(f"PARTIAL: Component 2 — Line-like shape found but position/stroke issues (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 — No gold horizontal line found on master slide")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Footer placeholder with 'Apex Dynamics' in white 10pt, bottom-left (0.25 points)
    try:
        footer_found = False
        footer_text_correct = False
        footer_font_correct = False

        for shape in master.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 3:
                # This is the footer placeholder (type FOOTER)
                footer_text = shape.text.strip()
                print(f"  Footer text: {repr(footer_text)}")

                if 'Apex Dynamics' in footer_text:
                    footer_text_correct = True

                # Check position - should be bottom-left (left < center of slide)
                if shape.left < slide_width / 2:
                    footer_found = True
                    print(f"  Footer position: left={shape.left}, top={shape.top} -- bottom-left")
                else:
                    print(f"  Footer position: left={shape.left} -- not bottom-left")

                # Check font properties: white color, 10pt
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                if run.font.color.type is not None:
                                    color = str(run.font.color.rgb).upper()
                                    if color == "FFFFFF":
                                        font_size = run.font.size
                                        # 10pt = 127000 EMU
                                        if font_size is not None and is_approximately_equal(font_size, 127000, tolerance=0.05):
                                            footer_font_correct = True
                                            print(f"  Footer font: white, {font_size} EMU ({font_size/12700:.1f}pt) -- correct")
                                        else:
                                            print(f"  Footer font: white but size={font_size} EMU, expected 127000 (10pt)")
                                    else:
                                        print(f"  Footer font color: #{color}, expected #FFFFFF")
                            except Exception:
                                pass
                break

        if footer_found and footer_text_correct and footer_font_correct:
            print(f"PASS: Component 3 — Footer 'Apex Dynamics' in white 10pt, bottom-left (0.25 pts)")
            total_score += 0.25
        elif footer_text_correct and footer_found:
            print(f"PARTIAL: Component 3 — Footer text correct, positioned left, but font issues (0.15 pts)")
            total_score += 0.15
        elif footer_text_correct:
            print(f"PARTIAL: Component 3 — Footer text correct but position/font wrong (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 3 — Footer placeholder missing or has wrong text")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide number placeholder positioned bottom-right (0.20 points)
    try:
        sldnum_found = False

        for shape in master.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 4:
                # This is the slide number placeholder
                # Should be positioned on the right half of the slide
                shape_right_edge = shape.left + shape.width
                slide_center = slide_width / 2

                print(f"  SlideNum position: left={shape.left}, right_edge={shape_right_edge}, slide_center={slide_center}")

                if shape.left > slide_center:
                    sldnum_found = True
                    print(f"  SlideNum is in the right half -- correct")
                else:
                    print(f"  SlideNum is NOT in the right half -- expected bottom-right")
                break

        if sldnum_found:
            print(f"PASS: Component 4 — Slide number placeholder bottom-right (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 4 — Slide number placeholder not positioned bottom-right")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
