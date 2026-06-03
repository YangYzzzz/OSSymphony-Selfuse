"""
Initial Setup: Create a 10-slide Apex Quarterly presentation with default white master
Task ID: impress_rp_002
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_002'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Apex Dynamics"
    slide1.placeholders[1].text = "Quarterly Business Review\nQ1 2026"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Financial Overview"
    for item in ["Product Development Update", "Market Expansion Strategy",
                 "Team & Hiring Pipeline", "Customer Success Metrics",
                 "Technology Roadmap", "Q2 Priorities"]:
        p = body2.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 3: Financial Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Financial Overview"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Revenue: $4.2M (+18% YoY)"
    for line in ["Gross Margin: 72.3%", "Operating Expenses: $2.8M",
                 "Net Income: $680K", "Cash Reserves: $12.1M",
                 "ARR Growth: 24% Quarter-over-Quarter"]:
        p = body3.add_paragraph()
        p.text = line
        p.level = 0

    # --- Slide 4: Product Development ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Product Development Update"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Platform v3.2 launched with AI-powered analytics"
    for line in ["Mobile app redesign completed (iOS & Android)",
                 "API response time improved by 40%",
                 "New integrations: Salesforce, HubSpot, Zendesk",
                 "99.97% uptime achieved in Q1"]:
        p = body4.add_paragraph()
        p.text = line
        p.level = 0

    # --- Slide 5: Market Expansion ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Market Expansion Strategy"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "EMEA office opened in London (March 2026)"
    for line in ["APAC partnership with Tanaka Solutions signed",
                 "Enterprise segment grew 31% QoQ",
                 "New verticals: Healthcare, FinTech",
                 "Channel partner program launching Q2"]:
        p = body5.add_paragraph()
        p.text = line
        p.level = 0

    # --- Slide 6: Team & Hiring ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Team & Hiring Pipeline"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Current headcount: 142 (+23 in Q1)"
    for line in ["Engineering: 58 | Sales: 34 | Operations: 28 | Marketing: 22",
                 "Key hires: VP of Engineering (Dr. Elena Vasquez)",
                 "Employee satisfaction score: 4.6/5.0",
                 "Retention rate: 94%"]:
        p = body6.add_paragraph()
        p.text = line
        p.level = 0

    # --- Slide 7: Customer Success ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Customer Success Metrics"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "NPS Score: 72 (Industry avg: 45)"
    for line in ["Active accounts: 1,247 (+180 new in Q1)",
                 "Churn rate: 2.1% (down from 3.8%)",
                 "Average contract value: $38,400",
                 "Support ticket resolution: 4.2 hours avg"]:
        p = body7.add_paragraph()
        p.text = line
        p.level = 0

    # --- Slide 8: Technology Roadmap ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Technology Roadmap"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Q2: Real-time collaboration engine"
    for line in ["Q3: Advanced reporting & custom dashboards",
                 "Q4: ML-based anomaly detection",
                 "2027 H1: Multi-tenant architecture upgrade",
                 "Continuous: Security hardening & SOC2 Type II"]:
        p = body8.add_paragraph()
        p.text = line
        p.level = 0

    # --- Slide 9: Q2 Priorities ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Q2 Priorities"
    body9 = slide9.placeholders[1].text_frame
    body9.text = "1. Close Series B funding round ($25M target)"
    for line in ["2. Launch EMEA go-to-market campaign",
                 "3. Ship collaboration engine (beta)",
                 "4. Expand enterprise sales team by 12 reps",
                 "5. Achieve SOC2 Type II certification"]:
        p = body9.add_paragraph()
        p.text = line
        p.level = 0

    # --- Slide 10: Thank You ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[0])
    slide10.shapes.title.text = "Thank You"
    slide10.placeholders[1].text = "Questions & Discussion\ncontact@apexdynamics.com"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
