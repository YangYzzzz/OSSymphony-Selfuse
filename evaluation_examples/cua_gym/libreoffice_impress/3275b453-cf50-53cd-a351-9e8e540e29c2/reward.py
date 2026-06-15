"""
Reward Script: Add a comparison table on slide 7 with programming languages
Task ID: impress_stu_045
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Table exists on slide 7 with 6 rows x 4 cols
  Component 2 (0.30): Header row content matches (Feature, Python, Java, C++)
  Component 3 (0.25): All 5 data rows have correct content
  Component 4 (0.20): Alternating row colors (#FFFFFF / #E8F0FE) on data rows
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_045'

# Expected table content
EXPECTED_HEADERS = ['Feature', 'Python', 'Java', 'C++']
EXPECTED_DATA = [
    ['Typing', 'Dynamic', 'Static', 'Static'],
    ['Speed', 'Moderate', 'Fast', 'Very Fast'],
    ['Learning Curve', 'Easy', 'Moderate', 'Hard'],
    ['Memory Management', 'Automatic', 'Automatic', 'Manual'],
    ['Primary Use', 'Data Science', 'Enterprise', 'Systems'],
]

# Expected alternating colors for data rows (rows 1-5)
# Odd data rows (1, 3, 5): FFFFFF; Even data rows (2, 4): E8F0FE
EXPECTED_ROW_COLORS = {
    1: 'FFFFFF',
    2: 'E8F0FE',
    3: 'FFFFFF',
    4: 'E8F0FE',
    5: 'FFFFFF',
}


def get_cell_fill_color(cell):
    """Extract the solid fill color from a table cell as hex string, or None."""
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is not None:
        solidFill = tcPr.find(qn('a:solidFill'))
        if solidFill is not None:
            srgb = solidFill.find(qn('a:srgbClr'))
            if srgb is not None:
                return srgb.get('val')
    return None


def find_table_on_slide(slide):
    """Find the first table shape on a slide, or None."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return shape.table
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Check presentation has at least 7 slides
    if len(prs.slides) < 7:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 7")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[6]  # Slide 7 (0-indexed)

    # Component 1: Table exists on slide 7 with correct dimensions (0.25 points)
    try:
        table = find_table_on_slide(slide)
        if table is not None:
            num_rows = len(table.rows)
            num_cols = len(table.columns)
            if num_rows == 6 and num_cols == 4:
                print(f"PASS: Component 1 - Table found on slide 7, {num_rows}x{num_cols} (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 - Table dimensions {num_rows}x{num_cols}, expected 6x4")
        else:
            print("FAIL: Component 1 - No table found on slide 7")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # If no valid table, remaining checks cannot proceed
    if table is None:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Header row content (0.30 points)
    try:
        if len(table.rows) >= 1 and len(table.columns) >= 4:
            actual_headers = [table.cell(0, c).text.strip() for c in range(4)]
            if actual_headers == EXPECTED_HEADERS:
                print(f"PASS: Component 2 - Header row matches: {actual_headers} (0.30 pts)")
                total_score += 0.30
            else:
                print(f"FAIL: Component 2 - Headers {actual_headers}, expected {EXPECTED_HEADERS}")
        else:
            print("FAIL: Component 2 - Table too small for header check")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Data rows content (0.25 points)
    try:
        if len(table.rows) >= 6 and len(table.columns) >= 4:
            correct_rows = 0
            for r_idx in range(5):
                actual_row = [table.cell(r_idx + 1, c).text.strip() for c in range(4)]
                expected_row = EXPECTED_DATA[r_idx]
                if actual_row == expected_row:
                    correct_rows += 1
                else:
                    print(f"  MISMATCH Row {r_idx + 1}: got {actual_row}, expected {expected_row}")

            if correct_rows == 5:
                print(f"PASS: Component 3 - All 5 data rows match (0.25 pts)")
                total_score += 0.25
            elif correct_rows >= 3:
                partial = round(0.25 * correct_rows / 5, 2)
                print(f"PARTIAL: Component 3 - {correct_rows}/5 rows correct ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 3 - Only {correct_rows}/5 data rows correct")
        else:
            print("FAIL: Component 3 - Table too small for data check")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Alternating row colors on data rows (0.20 points)
    try:
        if len(table.rows) >= 6 and len(table.columns) >= 4:
            correct_color_rows = 0
            for r_idx in range(1, 6):
                expected_color = EXPECTED_ROW_COLORS[r_idx].upper()
                # Check first cell of each data row as representative
                actual_color = get_cell_fill_color(table.cell(r_idx, 0))
                if actual_color is not None:
                    actual_color = actual_color.upper()

                if actual_color == expected_color:
                    correct_color_rows += 1
                else:
                    print(f"  COLOR MISMATCH Row {r_idx}: got {actual_color}, expected {expected_color}")

            if correct_color_rows == 5:
                print(f"PASS: Component 4 - All data rows have correct alternating colors (0.20 pts)")
                total_score += 0.20
            elif correct_color_rows >= 3:
                partial = round(0.20 * correct_color_rows / 5, 2)
                print(f"PARTIAL: Component 4 - {correct_color_rows}/5 rows with correct color ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 4 - Only {correct_color_rows}/5 data rows have correct colors")
        else:
            print("FAIL: Component 4 - Table too small for color check")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
