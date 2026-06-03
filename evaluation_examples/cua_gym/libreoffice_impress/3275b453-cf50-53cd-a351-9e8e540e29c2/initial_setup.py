"""
Initial Setup: Create a 9-slide CS_Intro presentation with slide 7 titled 'Language Comparison' but empty.
Task ID: impress_stu_045
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_045'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_body_slide(prs, title_text, body_lines, layout_idx=1):
    """Add a slide with title and bullet-point body content."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1].text_frame
    body.paragraphs[0].text = body_lines[0]
    for line in body_lines[1:]:
        p = body.add_paragraph()
        p.text = line
        p.level = 0
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Introduction to Computer Science"
    slide1.placeholders[1].text = "CS 101 - Fall 2025\nProfessor Elena Rodriguez"

    # --- Slide 2: Course Overview ---
    add_title_body_slide(prs, "Course Overview", [
        "Fundamentals of programming and computational thinking",
        "Data structures and algorithm design",
        "Object-oriented programming principles",
        "Software development lifecycle",
        "Introduction to multiple programming paradigms",
    ])

    # --- Slide 3: Data Types & Variables ---
    add_title_body_slide(prs, "Data Types & Variables", [
        "Integers: whole numbers (e.g., 42, -17, 0)",
        "Floating-point: decimal numbers (e.g., 3.14, -0.001)",
        "Strings: sequences of characters (e.g., \"Hello World\")",
        "Booleans: True or False values",
        "Arrays: ordered collections of elements",
        "Dictionaries: key-value pair mappings",
    ])

    # --- Slide 4: Control Flow ---
    add_title_body_slide(prs, "Control Flow Structures", [
        "Conditional statements: if, elif, else",
        "For loops: iterate over sequences",
        "While loops: repeat until condition is false",
        "Break and continue: loop control mechanisms",
        "Exception handling: try, except, finally blocks",
    ])

    # --- Slide 5: Functions & Modules ---
    add_title_body_slide(prs, "Functions & Modules", [
        "Defining reusable blocks of code with def/function",
        "Parameters, arguments, and return values",
        "Scope: local vs global variable visibility",
        "Lambda expressions for anonymous functions",
        "Importing and organizing code into modules",
    ])

    # --- Slide 6: Object-Oriented Programming ---
    add_title_body_slide(prs, "Object-Oriented Programming", [
        "Classes: blueprints for creating objects",
        "Encapsulation: bundling data and methods together",
        "Inheritance: extending functionality from parent classes",
        "Polymorphism: same interface, different behaviors",
        "Abstraction: hiding implementation complexity",
    ])

    # --- Slide 7: Language Comparison (EMPTY - title only) ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Language Comparison"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 8: Best Practices ---
    add_title_body_slide(prs, "Software Development Best Practices", [
        "Write clean, readable code with meaningful variable names",
        "Comment and document your code thoroughly",
        "Use version control systems like Git",
        "Test your code with unit tests and integration tests",
        "Follow the DRY principle: Don't Repeat Yourself",
    ])

    # --- Slide 9: Summary & Next Steps ---
    add_title_body_slide(prs, "Summary & Next Steps", [
        "Review key concepts from each module",
        "Complete lab assignments on algorithm design",
        "Start the final project: build a full application",
        "Office hours: Tuesdays and Thursdays 2-4 PM",
        "Course materials available on the student portal",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
