"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 34 needs a quick tweak: the first content text box should use 24 pt type and switch its font color to the preset “Dark Blue 2” (hex #002E6B). Can you apply that for me?
Generated: 2025-09-10 12:26:58
Status: success
Model: azure-o3
Total Steps: 2
"""

from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Pt
import os

# Path to the presentation provided in the task context
FILE_PATH = "/home/user/slide_34_needs_a_quick_tweak_the_first_content_text_box_should_use_24_pt_type_and_switch_its_font_co_golden.pptx"


def find_first_content_shape(slide):
    """Return the first text‐containing shape that is NOT a title placeholder."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                continue  # skip title placeholders
        return shape  # first qualifying content shape
    return None


def _extract_color(color_format):
    """Return (rgb_tuple, theme_color) from a ColorFormat object."""
    if color_format is None:
        return None, None
    if color_format.rgb is not None:
        rgb = color_format.rgb
        return (rgb[0], rgb[1], rgb[2]), None
    if color_format.theme_color is not None:
        return None, color_format.theme_color
    return None, None


def _effective_font(run, paragraph):
    """Derive effective size and color for a run, considering paragraph defaults."""
    size = run.font.size if run.font.size is not None else paragraph.font.size
    rgb, theme = _extract_color(run.font.color)
    if rgb is None and theme is None:
        rgb, theme = _extract_color(paragraph.font.color)
    return size, rgb, theme


def verify_font_properties(shape, expected_pt=24, expected_rgb=(0, 46, 107)):
    """Check that ALL runs in the shape use 24 pt and Dark Blue 2."""
    expected_size = Pt(expected_pt)
    tf = shape.text_frame

    total_runs = size_ok_runs = color_ok_runs = 0
    size_fail = color_fail = False

    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            total_runs += 1
            size, rgb, theme = _effective_font(run, paragraph)

            # ---- size check ----
            if size is None or abs(size - expected_size) >= 1000:  # 1000 EMU ≈ 0.08 pt tolerance
                size_fail = True
            else:
                size_ok_runs += 1

            # ---- color check ----
            if rgb is not None:
                if rgb == expected_rgb:
                    color_ok_runs += 1
                else:
                    color_fail = True
            elif theme is not None:
                if theme == MSO_THEME_COLOR.DARK_2:
                    color_ok_runs += 1
                else:
                    color_fail = True
            else:
                color_fail = True

    # Require every run to comply
    size_ok = not size_fail and total_runs == size_ok_runs
    color_ok = not color_fail and total_runs == color_ok_runs
    return size_ok, color_ok, total_runs, size_ok_runs, color_ok_runs


def verify_task(file_path):
    print(f"Verifying presentation at {file_path}")
    score = 0.0

    # 1) File existence & loading
    if not os.path.exists(file_path):
        print("✗ File not found.")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Unable to load PPTX: {e}")
        return 0.0

    # 2) Slide 34 presence
    target_idx = 33  # zero-based index
    if len(prs.slides) <= target_idx:
        print(f"✗ Slide 34 not found (presentation has {len(prs.slides)} slides).")
        return 0.0
    slide = prs.slides[target_idx]
    print("✓ Slide 34 located.")

    # 3) Locate first content shape
    shape = find_first_content_shape(slide)
    if shape is None:
        print("✗ No suitable content text box found on slide 34.")
        return 0.0
    print(f"✓ First content text box detected: '{shape.name}'")

    # 4) Verify font size & color
    size_ok, color_ok, total, size_runs, color_runs = verify_font_properties(shape)
    print(f"Runs inspected: {total}")
    print(f"24 pt runs:     {size_runs}/{total}")
    print(f"Dark Blue 2 runs: {color_runs}/{total}")

    if size_ok:
        print("✓ All text uses 24 pt (0.5 points)")
        score += 0.5
    else:
        print("✗ 24 pt size requirement not fully met")

    if color_ok:
        print("✓ All text uses Dark Blue 2 (0.5 points)")
        score += 0.5
    else:
        print("✗ Dark Blue 2 color requirement not fully met")

    final_score = min(score, 1.0)
    print(f"Final score: {final_score}")
    return final_score


if __name__ == "__main__":
    reward = verify_task(FILE_PATH)
    print(f"REWARD: {reward}")
