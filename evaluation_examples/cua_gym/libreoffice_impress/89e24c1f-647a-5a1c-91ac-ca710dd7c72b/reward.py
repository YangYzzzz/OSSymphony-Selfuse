"""
Reward Script: Science conference poster with title bar
Task ID: impress_ps_009
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.30): Slide dimensions 91.44 x 121.92 cm (portrait)
  - Component 2 (0.20): Title bar rectangle spanning full width at top
  - Component 3 (0.25): Title text correct content, 48pt bold white
  - Component 4 (0.25): Author text correct content, 28pt white
"""

import os

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_009'

# Expected dimensions in EMU: 91.44 cm x 121.92 cm
# 1 cm = 360000 EMU
EXPECTED_WIDTH_EMU = 32918400   # 91.44 cm
EXPECTED_HEIGHT_EMU = 43891200  # 121.92 cm
DIM_TOLERANCE = 0.02  # 2% tolerance

TITLE_TEXT = 'Neural Network Approaches to Climate Prediction'
AUTHOR_TEXT = 'Dr. A. Kumar, Dr. B. Zhang - MIT'


def approx_equal(a, b, tol=0.02):
    """Check approximate equality with relative tolerance."""
    if a == b:
        return True
    if a == 0 or b == 0:
        return abs(a - b) < 1000  # tiny absolute tolerance for zero
    return abs(a - b) / max(abs(a), abs(b)) <= tol


def persist_app_state(domain):
    """Save any unsaved LibreOffice edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    if len(prs.slides) == 0:
        print("CRITICAL: No slides found")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]

    # Component 1: Slide dimensions are 91.44 x 121.92 cm portrait (0.30 points)
    try:
        w = prs.slide_width
        h = prs.slide_height
        width_ok = approx_equal(w, EXPECTED_WIDTH_EMU)
        height_ok = approx_equal(h, EXPECTED_HEIGHT_EMU)
        portrait_ok = h > w

        if width_ok and height_ok and portrait_ok:
            print(f"PASS: Component 1 -- Slide dimensions correct: {w}x{h} EMU (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 1 -- Expected ~{EXPECTED_WIDTH_EMU}x{EXPECTED_HEIGHT_EMU} EMU portrait, "
                  f"got {w}x{h} (width_ok={width_ok}, height_ok={height_ok}, portrait={portrait_ok})")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Title bar rectangle spanning full width at top (0.20 points)
    try:
        rect_found = False
        for shape in slide.shapes:
            # Look for a rectangle/auto shape that spans full width and is at the top
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                is_full_width = approx_equal(shape.width, prs.slide_width, tol=0.05)
                is_at_top = shape.top < prs.slide_height * 0.1  # top 10%
                is_bar_height = shape.height > 0 and shape.height < prs.slide_height * 0.25
                if is_full_width and is_at_top and is_bar_height:
                    rect_found = True
                    break

        if rect_found:
            print(f"PASS: Component 2 -- Title bar rectangle found at top, full width (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 2 -- No full-width rectangle found at top of slide")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Title text with correct content, ~48pt bold white (0.25 points)
    try:
        title_found = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            full_text = shape.text_frame.text.strip()
            if TITLE_TEXT.lower() in full_text.lower():
                # Check font properties on runs
                sub_score = 0.0
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if TITLE_TEXT.lower() in run.text.strip().lower():
                            # Text match
                            sub_score += 0.10

                            # Bold check
                            if run.font.bold:
                                sub_score += 0.05
                            else:
                                print(f"  INFO: Title run not bold (bold={run.font.bold})")

                            # Size check (~48pt = 609600 EMU, allow some tolerance)
                            if run.font.size is not None and approx_equal(run.font.size, Pt(48), tol=0.10):
                                sub_score += 0.05
                            else:
                                print(f"  INFO: Title font size {run.font.size}, expected ~{Pt(48)}")

                            # White color check
                            try:
                                if run.font.color.type is not None:
                                    rgb_str = str(run.font.color.rgb).upper()
                                    if rgb_str == 'FFFFFF':
                                        sub_score += 0.05
                                    else:
                                        print(f"  INFO: Title color {rgb_str}, expected FFFFFF")
                                else:
                                    print(f"  INFO: Title color type is None")
                            except Exception:
                                print(f"  INFO: Could not read title color")

                            title_found = True
                            break
                    if title_found:
                        break

                if title_found:
                    total_score += sub_score
                    print(f"PASS: Component 3 -- Title text found with formatting ({sub_score} pts)")
                break

        if not title_found:
            print(f"FAIL: Component 3 -- Title text '{TITLE_TEXT}' not found")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Author text with correct content, ~28pt white (0.25 points)
    try:
        author_found = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            full_text = shape.text_frame.text.strip()
            if AUTHOR_TEXT.lower() in full_text.lower():
                sub_score = 0.0
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if AUTHOR_TEXT.lower() in run.text.strip().lower():
                            # Text match
                            sub_score += 0.10

                            # Size check (~28pt = 355600 EMU)
                            if run.font.size is not None and approx_equal(run.font.size, Pt(28), tol=0.15):
                                sub_score += 0.075
                            else:
                                print(f"  INFO: Author font size {run.font.size}, expected ~{Pt(28)}")

                            # White color check
                            try:
                                if run.font.color.type is not None:
                                    rgb_str = str(run.font.color.rgb).upper()
                                    if rgb_str == 'FFFFFF':
                                        sub_score += 0.075
                                    else:
                                        print(f"  INFO: Author color {rgb_str}, expected FFFFFF")
                                else:
                                    print(f"  INFO: Author color type is None")
                            except Exception:
                                print(f"  INFO: Could not read author color")

                            author_found = True
                            break
                    if author_found:
                        break

                if author_found:
                    total_score += sub_score
                    print(f"PASS: Component 4 -- Author text found with formatting ({sub_score} pts)")
                break

        if not author_found:
            print(f"FAIL: Component 4 -- Author text '{AUTHOR_TEXT}' not found")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist app state before verification
persist_app_state('libreoffice_impress')

# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
