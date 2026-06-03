"""
Initial Setup: Create an empty presentation for a science conference poster task.
Task ID: impress_ps_009
Domain: libreoffice_impress

The initial file has 1 empty slide at default dimensions (10 x 7.5 inches landscape).
No custom sizes, no shapes, no text.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_009'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Default dimensions: 10 x 7.5 inches (25.4 x 19.05 cm) - landscape
    # Add one blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only layout - blank-ish

    # Remove any default placeholder shapes to ensure the slide is completely empty
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress for GUI-ready state
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
