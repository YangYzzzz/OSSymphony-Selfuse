"""
Reward Script: Move comparison table to bottom of slide 3
Task ID: osworld_impress_table_position_bottom_004
Domain: libreoffice_impress
Scoring:
  Component 1 (0.6 pts): Table top is in lower portion of slide (top > 55% of slide height)
  Component 2 (0.4 pts): Table bottom is near the bottom of slide (> 95% of slide height)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_table_position_bottom_004'


def persist_app_state():
    """Best-effort save in case LibreOffice has unsaved GUI edits."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify that the comparison table on slide 3 has been repositioned to the
    lower/bottom portion of the slide.

    Scoring rubric:
      Component 1 (0.6 pts): Table top is in the lower portion (> 55% slide height).
          Initial state: table_top = 2,926,080 EMU ~ 42.7% of slide — FAILS this check.
          Golden state:  table_top = 4,480,560 EMU ~ 65.3% of slide — PASSES this check.

      Component 2 (0.4 pts): Table bottom is near the bottom edge (> 95% slide height).
          Initial state: table_bottom = 5,212,080 EMU ~ 76.0% of slide — FAILS this check.
          Golden state:  table_bottom = 6,766,560 EMU ~ 98.7% of slide — PASSES this check.

    Returns float in [0.0, 1.0].
    """
    total_score = 0.0

    # --- Load presentation ---
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Validate we have at least 3 slides
    if len(prs.slides) < 3:
        print(f"CRITICAL: Expected at least 3 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide_height = prs.slide_height  # 6858000 EMU for standard widescreen
    slide = prs.slides[2]  # Slide 3 is index 2

    # Find the table shape on slide 3
    table_shape = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_shape = shape
            break

    if table_shape is None:
        print("CRITICAL: No table found on slide 3")
        print("REWARD: 0.0")
        return 0.0

    table_top = table_shape.top
    table_height_emu = table_shape.height
    table_bottom = table_top + table_height_emu

    top_ratio = table_top / slide_height
    bottom_ratio = table_bottom / slide_height

    print(f"INFO: slide_height = {slide_height} EMU")
    print(f"INFO: table_top    = {table_top} EMU  ({top_ratio:.3f} of slide height)")
    print(f"INFO: table_bottom = {table_bottom} EMU  ({bottom_ratio:.3f} of slide height)")

    # --- Component 1: Table top is in lower portion of slide (> 55% of slide height) ---
    # Threshold: 55% of 6858000 = 3771900 EMU
    # This distinguishes middle (42.7%) from lower portion (65.3%)
    LOWER_PORTION_THRESHOLD = 0.55  # 55% of slide height
    lower_threshold_emu = slide_height * LOWER_PORTION_THRESHOLD
    try:
        if table_top > lower_threshold_emu:
            print(f"PASS: Component 1 — Table top ({table_top} EMU = {top_ratio:.1%}) is in lower portion "
                  f"(> {lower_threshold_emu:.0f} EMU = {LOWER_PORTION_THRESHOLD:.0%} threshold) (0.6 pts)")
            total_score += 0.6
        else:
            print(f"FAIL: Component 1 — Table top ({table_top} EMU = {top_ratio:.1%}) is NOT in lower portion "
                  f"(required > {lower_threshold_emu:.0f} EMU = {LOWER_PORTION_THRESHOLD:.0%})")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # --- Component 2: Table bottom is near the bottom of the slide (> 95% of slide height) ---
    # Threshold: 95% of 6858000 = 6515100 EMU
    # Golden: 6766560 (98.7%) passes. Initial: 5212080 (76.0%) fails.
    NEAR_BOTTOM_THRESHOLD = 0.95  # 95% of slide height
    near_bottom_emu = slide_height * NEAR_BOTTOM_THRESHOLD
    try:
        if table_bottom > near_bottom_emu:
            print(f"PASS: Component 2 — Table bottom ({table_bottom} EMU = {bottom_ratio:.1%}) is near slide bottom "
                  f"(> {near_bottom_emu:.0f} EMU = {NEAR_BOTTOM_THRESHOLD:.0%} threshold) (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 2 — Table bottom ({table_bottom} EMU = {bottom_ratio:.1%}) is NOT near slide bottom "
                  f"(required > {near_bottom_emu:.0f} EMU = {NEAR_BOTTOM_THRESHOLD:.0%})")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# --- Entrypoint ---
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
