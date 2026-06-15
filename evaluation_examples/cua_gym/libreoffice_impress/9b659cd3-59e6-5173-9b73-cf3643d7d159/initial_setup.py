"""
Initial Setup: Product comparison presentation with table floating in middle of slide 3
Task ID: osworld_impress_table_position_bottom_004
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_table_position_bottom_004'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard widescreen 10 x 7.5 inches
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_layouts = prs.slide_layouts
    # layout 0 = Title Slide, 1 = Title+Content, 5 = Blank, 6 = Title Only

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(slide_layouts[0])
    slide1.shapes.title.text = "Product Comparison 2025"
    slide1.placeholders[1].text = "Q2 Competitive Analysis\nMarketing Strategy Division"

    # ---- Slide 2: Executive Summary ----
    slide2 = prs.slides.add_slide(slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Key findings from our Q2 product benchmarking study"
    bullets2 = [
        "Market share analysis across 5 product categories",
        "Price-performance ratios evaluated for 12 competitors",
        "Customer satisfaction scores from 2,400 survey responses",
        "Feature gap analysis completed for top 3 segments",
    ]
    for b in bullets2:
        p = tf2.add_paragraph()
        p.text = b
        p.level = 1

    # ---- Slide 3: Feature Comparison (with table in MIDDLE - initial state) ----
    slide3 = prs.slides.add_slide(slide_layouts[5])  # Blank layout

    # Title
    title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.8))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "Feature Comparison: Our Product vs Competitors"
    run_title = p_title.runs[0]
    run_title.font.size = Pt(28)
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Bullet points (3 bullets)
    bullet_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(8.5), Inches(1.8))
    tf_bullets = bullet_box.text_frame
    tf_bullets.word_wrap = True

    bullet_items = [
        "NovaPro X leads in AI-assisted workflow automation with 94% task accuracy",
        "Competitor pricing averages 18% higher for equivalent feature sets",
        "Integration ecosystem: NovaPro X supports 340+ third-party connectors",
    ]
    for i, item in enumerate(bullet_items):
        if i == 0:
            p = tf_bullets.paragraphs[0]
        else:
            p = tf_bullets.add_paragraph()
        p.text = f"\u2022 {item}"
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Comparison table — positioned in the MIDDLE of the slide (initial state, NOT at bottom)
    # Table is at vertical center (~3.2 inches top), which is in the middle of the slide
    table_left = Inches(0.5)
    table_top = Inches(3.2)   # Middle of slide (initial: floating in middle)
    table_width = Inches(9.0)
    table_height = Inches(2.5)

    rows = 5
    cols = 4
    table_shape = slide3.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
    table = table_shape.table

    # Set column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(2.0)
    table.columns[3].width = Inches(2.5)

    # Header row
    headers = ["Feature", "NovaPro X", "CompetitorA", "CompetitorB"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.size = Pt(13)
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # Header background
        from pptx.oxml.ns import qn
        from lxml import etree
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', '1F497D')

    # Data rows
    data_rows = [
        ["AI Automation",       "94% accuracy",   "78% accuracy",   "81% accuracy"],
        ["Cloud Integration",   "340+ connectors", "180+ connectors", "220+ connectors"],
        ["Security Compliance", "SOC2, GDPR, ISO", "SOC2, GDPR",     "GDPR, ISO"],
        ["Pricing (per seat)",  "$49/month",       "$58/month",      "$55/month"],
    ]
    for row_idx, row_data in enumerate(data_rows, 1):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(12)
                    # Alternate row shading
                    if row_idx % 2 == 0:
                        from pptx.oxml.ns import qn
                        from lxml import etree
                        tc = cell._tc
                        tcPr = tc.get_or_add_tcPr()
                        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
                        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
                        srgbClr.set('val', 'DCE6F1')

    # ---- Slide 4: Pricing Analysis ----
    slide4 = prs.slides.add_slide(slide_layouts[1])
    slide4.shapes.title.text = "Pricing Analysis"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Competitive pricing landscape for enterprise tier"
    pricing_bullets = [
        "NovaPro X: $49/seat/month — 15% below market average",
        "CompetitorA: $58/seat/month — includes limited support",
        "CompetitorB: $55/seat/month — requires annual commitment",
        "Volume discounts available for 50+ seat deployments",
    ]
    for b in pricing_bullets:
        p = tf4.add_paragraph()
        p.text = b
        p.level = 1

    # ---- Slide 5: Customer Satisfaction ----
    slide5 = prs.slides.add_slide(slide_layouts[1])
    slide5.shapes.title.text = "Customer Satisfaction Scores"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Net Promoter Score (NPS) benchmarking — Q2 2025"
    csat_bullets = [
        "NovaPro X NPS: 72 (industry average: 45)",
        "Customer retention rate: 94% year-over-year",
        "Support ticket resolution: avg 4.2 hours",
        "Feature request implementation: 68% fulfilled within 2 quarters",
    ]
    for b in csat_bullets:
        p = tf5.add_paragraph()
        p.text = b
        p.level = 1

    # ---- Slide 6: Conclusion & Recommendations ----
    slide6 = prs.slides.add_slide(slide_layouts[1])
    slide6.shapes.title.text = "Conclusion & Next Steps"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Strategic recommendations based on competitive analysis"
    conclusion_bullets = [
        "Maintain pricing advantage — do not exceed $52/seat in next cycle",
        "Invest in connector ecosystem to reach 400+ by Q4 2025",
        "Leverage NPS leadership in enterprise sales collateral",
        "Initiate partnership talks with top 3 integration vendors",
    ]
    for b in conclusion_bullets:
        p = tf6.add_paragraph()
        p.text = b
        p.level = 1

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open initial file in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
