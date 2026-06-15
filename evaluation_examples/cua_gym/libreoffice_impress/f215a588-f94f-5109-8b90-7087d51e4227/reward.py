"""
Reward Script: World map slide with colored pins, city labels, and legend
Task ID: impress_sales_065
Domain: libreoffice_impress
Scoring:
  Component 1: World map image on slide 6 (0.2 pts)
  Component 2: 5 colored pin marker shapes (0.3 pts)
  Component 3: 5 city labels with employee counts (0.3 pts)
  Component 4: Legend in bottom-right corner (0.2 pts)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_065'

# Expected cities with colors and employee counts
EXPECTED_PINS = {
    'New York':  {'color': 'FF0000', 'employees': '150'},
    'London':    {'color': '0000FF', 'employees': '120'},
    'Tokyo':     {'color': '00AA00', 'employees': '85'},
    'Sydney':    {'color': 'FF6B35', 'employees': '65'},
    'São Paulo': {'color': '9C27B0', 'employees': '45'},
}

CITY_NAMES = list(EXPECTED_PINS.keys())


def get_shape_fill_rgb(shape):
    """Get the fill color RGB string of a shape, or None if not solid fill."""
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[5]  # Slide 6 (0-indexed)
    shapes = list(slide.shapes)

    # Component 1: World map image on slide 6 (0.2 points)
    # The initial slide has NO images. The golden slide has a world_map.png image.
    try:
        pictures = [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        if len(pictures) >= 1:
            # Verify the image is reasonably large (map should be substantial)
            pic = pictures[0]
            pic_area = pic.width * pic.height
            slide_area = prs.slide_width * prs.slide_height
            if pic_area > slide_area * 0.1:  # image covers >10% of slide
                print(f"PASS: Component 1 — World map image found, size ({pic.width}x{pic.height}) (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 1 — Image found but too small: {pic.width}x{pic.height}")
        else:
            print(f"FAIL: Component 1 — No image found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: 5 colored pin marker shapes (0.3 points, 0.06 each)
    # Pin markers are oval/auto shapes with specific fill colors.
    # The initial slide has NO auto shapes. We look for ovals with the expected colors.
    try:
        auto_shapes = [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
        # Separate large pins (marker pins on map) from small pins (legend dots)
        # Map pins are larger (274320 EMU ~0.3 inch), legend dots are smaller (109728 EMU ~0.12 inch)
        large_ovals = [s for s in auto_shapes if s.width > 150000 and s.width < 400000
                       and s.height > 150000 and s.height < 400000]

        found_colors = set()
        for oval in large_ovals:
            rgb = get_shape_fill_rgb(oval)
            if rgb:
                found_colors.add(rgb)

        pin_score = 0.0
        for city, info in EXPECTED_PINS.items():
            expected_color = info['color']
            if expected_color in found_colors:
                print(f"  PASS: Pin for {city} — color {expected_color} found")
                pin_score += 0.06
            else:
                print(f"  FAIL: Pin for {city} — color {expected_color} not found among large ovals")

        if pin_score > 0:
            print(f"PASS: Component 2 — {pin_score/0.06:.0f}/5 pin markers found ({pin_score:.2f} pts)")
        else:
            print(f"FAIL: Component 2 — No pin markers with expected colors found")
        total_score += pin_score
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: 5 city labels with employee counts (0.3 points, 0.06 each)
    # Each label should contain city name and employee count like "New York - 150 employees"
    # The initial slide has NO such text boxes.
    try:
        text_shapes = [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
                       or (hasattr(s, 'text_frame') and s.has_text_frame)]
        all_texts = []
        for s in text_shapes:
            try:
                txt = s.text.strip() if hasattr(s, 'text') and s.text else ''
                if txt:
                    all_texts.append(txt)
            except Exception:
                pass

        label_score = 0.0
        for city, info in EXPECTED_PINS.items():
            emp_count = info['employees']
            # Look for a text that contains both the city name and employee count
            found_label = False
            for txt in all_texts:
                if city.lower() in txt.lower() and emp_count in txt:
                    found_label = True
                    break
            if found_label:
                print(f"  PASS: Label for {city} with {emp_count} employees found")
                label_score += 0.06
            else:
                print(f"  FAIL: Label for {city} with {emp_count} employees not found")

        if label_score > 0:
            print(f"PASS: Component 3 — {label_score/0.06:.0f}/5 city labels found ({label_score:.2f} pts)")
        else:
            print(f"FAIL: Component 3 — No city labels with employee counts found")
        total_score += label_score
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Legend in bottom-right corner (0.2 points)
    # Legend should be in bottom-right area of slide, containing city names.
    # We look for a group of elements in the bottom-right quadrant that form a legend.
    # The initial slide has NO legend.
    try:
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        # Look for a container shape (like rounded rectangle) in bottom-right
        # Bottom-right means left > 50% of slide_w AND top > 50% of slide_h
        legend_found = False
        legend_has_title = False
        legend_city_count = 0

        # Check for rounded rectangle or similar container in bottom-right
        container_shapes = [s for s in auto_shapes
                            if s.left > slide_w * 0.5 and s.top > slide_h * 0.5
                            and s.width > 500000]  # reasonably large
        # Also check for text boxes in bottom-right that might be legend entries
        br_text_shapes = [s for s in shapes
                          if hasattr(s, 'text') and s.text
                          and s.left > slide_w * 0.5 and s.top > slide_h * 0.5]

        if len(container_shapes) >= 1:
            legend_found = True
            print(f"  Legend container found at ({container_shapes[0].left}, {container_shapes[0].top})")

        # Check for legend title text (like "Office Locations" or "Legend")
        for s in br_text_shapes:
            txt = s.text.strip().lower()
            if any(kw in txt for kw in ['legend', 'office', 'location', 'offices']):
                legend_has_title = True

        # Check for city names in legend area (small text items in bottom-right)
        for s in br_text_shapes:
            txt = s.text.strip()
            for city in CITY_NAMES:
                if city.lower() in txt.lower() and len(txt) < 30:  # short label, not the main city label
                    legend_city_count += 1
                    break

        # Score: legend container + title + at least 3 cities referenced
        legend_score = 0.0
        if legend_found and legend_has_title and legend_city_count >= 3:
            legend_score = 0.2
            print(f"PASS: Component 4 — Full legend found with title and {legend_city_count} cities (0.2 pts)")
        elif legend_found and legend_city_count >= 3:
            legend_score = 0.15
            print(f"PARTIAL: Component 4 — Legend found with {legend_city_count} cities but no clear title (0.15 pts)")
        elif legend_found:
            legend_score = 0.1
            print(f"PARTIAL: Component 4 — Legend container found but missing city entries (0.1 pts)")
        else:
            print(f"FAIL: Component 4 — No legend found in bottom-right corner")

        total_score += legend_score
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
