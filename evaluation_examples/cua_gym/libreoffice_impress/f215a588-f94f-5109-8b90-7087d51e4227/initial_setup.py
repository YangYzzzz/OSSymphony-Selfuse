"""
Initial Setup: Create a 9-slide Global Pitch presentation with slide 6 titled 'Global Presence'
              and a world_map.png file on the Desktop.
Task ID: impress_sales_065
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_065'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
WORLD_MAP = f'{WORKDIR}/Desktop/world_map.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_world_map_image():
    """Create a simple world map placeholder image on the Desktop."""
    os.makedirs(f'{WORKDIR}/Desktop', exist_ok=True)
    img = Image.new('RGB', (1200, 600), color=(200, 220, 240))
    draw = ImageDraw.Draw(img)

    # Draw simplified continent shapes as filled polygons
    # North America
    na_points = [(120, 80), (300, 60), (340, 120), (320, 200), (280, 260),
                 (200, 280), (140, 240), (100, 180), (90, 120)]
    draw.polygon(na_points, fill=(140, 180, 120), outline=(80, 120, 80))

    # South America
    sa_points = [(250, 300), (310, 290), (340, 340), (330, 420), (300, 480),
                 (260, 500), (230, 460), (220, 380), (230, 330)]
    draw.polygon(sa_points, fill=(140, 180, 120), outline=(80, 120, 80))

    # Europe
    eu_points = [(520, 70), (600, 60), (640, 90), (630, 140), (600, 170),
                 (560, 160), (530, 130), (510, 100)]
    draw.polygon(eu_points, fill=(140, 180, 120), outline=(80, 120, 80))

    # Africa
    af_points = [(520, 180), (600, 170), (640, 220), (630, 340), (590, 400),
                 (550, 390), (510, 320), (500, 240)]
    draw.polygon(af_points, fill=(140, 180, 120), outline=(80, 120, 80))

    # Asia
    asia_points = [(650, 50), (850, 60), (950, 100), (980, 180), (940, 260),
                   (870, 280), (780, 260), (720, 200), (680, 150), (640, 100)]
    draw.polygon(asia_points, fill=(140, 180, 120), outline=(80, 120, 80))

    # Australia
    au_points = [(880, 360), (960, 340), (1020, 370), (1030, 420), (990, 460),
                 (930, 450), (880, 410)]
    draw.polygon(au_points, fill=(140, 180, 120), outline=(80, 120, 80))

    # Add grid lines for map feel
    for x in range(0, 1200, 120):
        draw.line([(x, 0), (x, 600)], fill=(180, 200, 220), width=1)
    for y in range(0, 600, 60):
        draw.line([(0, y), (1200, y)], fill=(180, 200, 220), width=1)

    img.save(WORLD_MAP)
    print(f'World map image created: {WORLD_MAP}')


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a textbox with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "TechVista Global Solutions"
    slide1.placeholders[1].text = "Strategic Growth & Expansion Plan 2026"
    bg1 = slide1.background.fill
    bg1.solid()
    bg1.fore_color.rgb = RGBColor(0x0A, 0x2A, 0x4A)
    for shape in slide1.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # ---- Slide 2: Company Overview ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide2, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                "Company Overview", font_size=32, bold=True,
                color=RGBColor(0x0A, 0x2A, 0x4A))
    overview_text = (
        "Founded in 2012, TechVista Global Solutions has grown from a 15-person "
        "startup in San Francisco to a multinational enterprise with over 2,400 "
        "employees across 5 continents. Our core offerings include enterprise cloud "
        "migration, AI-driven analytics platforms, and cybersecurity consulting. "
        "In FY2025, we achieved $487M in revenue with 23% year-over-year growth."
    )
    add_textbox(slide2, Inches(0.5), Inches(1.5), Inches(12), Inches(4),
                overview_text, font_size=16)

    # ---- Slide 3: Revenue Breakdown ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide3, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                "Revenue Breakdown by Segment", font_size=32, bold=True,
                color=RGBColor(0x0A, 0x2A, 0x4A))
    rows, cols = 6, 3
    tbl_shape = slide3.shapes.add_table(rows, cols, Inches(1), Inches(1.8),
                                        Inches(10), Inches(3.5))
    tbl = tbl_shape.table
    headers = ["Business Segment", "FY2025 Revenue", "Growth %"]
    data = [
        ["Cloud Migration Services", "$178.2M", "28%"],
        ["AI Analytics Platform", "$134.5M", "41%"],
        ["Cybersecurity Consulting", "$98.7M", "15%"],
        ["Managed IT Services", "$52.3M", "12%"],
        ["Training & Certification", "$23.3M", "8%"],
    ]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            tbl.cell(r, c).text = val

    # ---- Slide 4: Key Clients ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide4, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                "Key Enterprise Clients", font_size=32, bold=True,
                color=RGBColor(0x0A, 0x2A, 0x4A))
    clients_text = (
        "Fortune 500 Partners:\n"
        "  - Meridian Financial Group (since 2018) - Full cloud migration\n"
        "  - Atlas Healthcare Network (since 2019) - AI diagnostics platform\n"
        "  - NovaTech Manufacturing (since 2020) - IoT analytics suite\n"
        "  - Pinnacle Retail Holdings (since 2021) - Cybersecurity overhaul\n"
        "  - Horizon Energy Corp (since 2022) - Predictive maintenance AI\n\n"
        "Government & Public Sector:\n"
        "  - Singapore Digital Infrastructure Board\n"
        "  - UK National Health Analytics Project\n"
        "  - City of Toronto Smart Transit Initiative"
    )
    add_textbox(slide4, Inches(0.5), Inches(1.5), Inches(11), Inches(5),
                clients_text, font_size=14)

    # ---- Slide 5: Growth Strategy ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide5, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                "Growth Strategy 2026-2028", font_size=32, bold=True,
                color=RGBColor(0x0A, 0x2A, 0x4A))
    strategy_text = (
        "1. Market Expansion: Enter Southeast Asian and Middle Eastern markets\n"
        "2. Product Innovation: Launch TechVista Nexus - unified enterprise platform\n"
        "3. Strategic Acquisitions: Target 2-3 cybersecurity firms in EMEA region\n"
        "4. Talent Growth: Expand workforce to 3,500+ by end of 2027\n"
        "5. Sustainability: Achieve carbon-neutral operations by 2028"
    )
    add_textbox(slide5, Inches(0.5), Inches(1.5), Inches(11), Inches(5),
                strategy_text, font_size=16)

    # ---- Slide 6: Global Presence (empty content - task target) ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide6, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                "Global Presence", font_size=32, bold=True,
                color=RGBColor(0x0A, 0x2A, 0x4A))
    # Content area intentionally left empty for the task

    # ---- Slide 7: Team Leadership ----
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide7, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                "Executive Leadership", font_size=32, bold=True,
                color=RGBColor(0x0A, 0x2A, 0x4A))
    team_text = (
        "CEO: Dr. Elena Vasquez - Former VP of Engineering at Microsoft Azure\n"
        "CTO: Raj Patel - Co-founder, 20+ years in distributed systems\n"
        "CFO: Margaret O'Brien - Previously at Goldman Sachs Technology Division\n"
        "COO: James Nakamura - Ex-McKinsey, specializing in tech operations\n"
        "VP Sales: Sofia Bergstrom - Built $200M pipeline at Salesforce"
    )
    add_textbox(slide7, Inches(0.5), Inches(1.5), Inches(11), Inches(5),
                team_text, font_size=15)

    # ---- Slide 8: Financial Projections ----
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide8, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                "Financial Projections", font_size=32, bold=True,
                color=RGBColor(0x0A, 0x2A, 0x4A))
    rows2, cols2 = 5, 4
    tbl2_shape = slide8.shapes.add_table(rows2, cols2, Inches(1), Inches(1.8),
                                         Inches(10), Inches(3))
    tbl2 = tbl2_shape.table
    headers2 = ["Metric", "FY2026 (Proj)", "FY2027 (Proj)", "FY2028 (Proj)"]
    data2 = [
        ["Total Revenue", "$612M", "$780M", "$985M"],
        ["EBITDA Margin", "22%", "25%", "28%"],
        ["Headcount", "2,900", "3,500", "4,200"],
        ["Office Locations", "8", "11", "14"],
    ]
    for c, h in enumerate(headers2):
        cell = tbl2.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row_data in enumerate(data2, 1):
        for c, val in enumerate(row_data):
            tbl2.cell(r, c).text = val

    # ---- Slide 9: Thank You / Contact ----
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    bg9 = slide9.background.fill
    bg9.solid()
    bg9.fore_color.rgb = RGBColor(0x0A, 0x2A, 0x4A)
    add_textbox(slide9, Inches(2), Inches(2), Inches(9), Inches(1.5),
                "Thank You", font_size=44, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_textbox(slide9, Inches(2), Inches(3.8), Inches(9), Inches(2),
                "Contact: invest@techvista.com | +1 (415) 555-0192\nwww.techvista-global.com",
                font_size=18, color=RGBColor(0xCC, 0xCC, 0xCC),
                alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_world_map_image()
create_initial()
