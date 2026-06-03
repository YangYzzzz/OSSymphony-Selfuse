"""
Reward Script: Create award/recognition slide at position 9 with 4 star badges and title
Task ID: impress_sales_084
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Title text "Award-Winning Platform" exists on slide 9 in 32pt bold
  Component 2 (0.35): Four 5-point star shapes with gold (#FFD700) fill on slide 9
  Component 3 (0.25): Four award label text boxes with correct names below stars
  Component 4 (0.15): Stars are arranged in a single row (similar top positions) and evenly spaced
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_084'

EXPECTED_AWARDS = ['Best SaaS 2024', 'G2 Leader', 'Gartner Cool Vendor', 'Forbes Cloud 100']


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 9 slides
    if len(prs.slides) < 9:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 9")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[8]  # Slide 9 (0-indexed)

    # Component 1: Title text "Award-Winning Platform" in 32pt bold (0.25 points)
    # This checks for a NEW text element added by the task (not the empty placeholder)
    try:
        title_flags = {'found': 0, 'bold': 0, 'size': 0}  # 0=no, 1=yes

        for shape in slide.shapes:
            if not hasattr(shape, 'text_frame'):
                continue
            full_text = shape.text_frame.text.strip() if shape.text_frame.text else ""
            if 'award-winning platform' in full_text.lower():
                title_flags['found'] = 1
                # Check font properties on runs
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if 'award-winning platform' in run.text.lower():
                            # Check bold (None means inherit, we accept True)
                            if run.font.bold is True:
                                title_flags['bold'] = 1
                            # Check size: 32pt = 406400 EMU
                            if run.font.size is not None:
                                # Allow some tolerance: 30-34pt
                                size_pt = run.font.size / 12700  # EMU to pt
                                if 30 <= size_pt <= 34:
                                    title_flags['size'] = 1

        if title_flags['found'] and title_flags['bold'] and title_flags['size']:
            print(f"PASS: Component 1 — Title 'Award-Winning Platform' found, 32pt bold (0.25 pts)")
            total_score += 0.25
        elif title_flags['found'] and (title_flags['bold'] or title_flags['size']):
            print(f"PARTIAL: Component 1 — Title found but missing {'bold' if not title_flags['bold'] else '32pt size'} (0.1 pts)")
            total_score += 0.1
        elif title_flags['found']:
            print(f"PARTIAL: Component 1 — Title text found but not bold and not 32pt (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 1 — Title 'Award-Winning Platform' not found on slide 9")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Four 5-point star shapes with gold (#FFD700) fill (0.35 points)
    try:
        star_shapes = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                try:
                    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
                    if shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.STAR_5_POINT:
                        star_shapes.append(shape)
                except:
                    # Fallback: check shape name contains "Star"
                    if 'star' in shape.name.lower():
                        star_shapes.append(shape)

        star_count = len(star_shapes)
        gold_stars = 0

        for star in star_shapes:
            try:
                fill = star.fill
                if fill.type is not None and fill.type == 1:  # SOLID
                    rgb = str(fill.fore_color.rgb).upper()
                    if rgb == 'FFD700':
                        gold_stars += 1
            except:
                pass

        if star_count >= 4 and gold_stars >= 4:
            print(f"PASS: Component 2 — {star_count} star shapes found, {gold_stars} with gold fill (0.35 pts)")
            total_score += 0.35
        elif star_count >= 4 and gold_stars >= 2:
            print(f"PARTIAL: Component 2 — {star_count} stars but only {gold_stars} gold (0.2 pts)")
            total_score += 0.2
        elif star_count >= 2:
            print(f"PARTIAL: Component 2 — Only {star_count} star shapes found (need 4) (0.1 pts)")
            total_score += 0.1
        else:
            print(f"FAIL: Component 2 — Found {star_count} star shapes (need 4)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Four award label text boxes with correct names (0.25 points)
    try:
        # Collect all non-title text shapes on slide 9 that are text boxes
        label_texts = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                text = shape.text_frame.text.strip() if hasattr(shape, 'text_frame') else ""
                if text and 'award-winning platform' not in text.lower():
                    label_texts.append(text)

        found_awards = []
        for expected in EXPECTED_AWARDS:
            for label in label_texts:
                if expected.lower() in label.lower():
                    found_awards.append(expected)
                    break

        match_count = len(found_awards)
        points_per_label = 0.25 / 4  # ~0.0625 per label

        if match_count == 4:
            print(f"PASS: Component 3 — All 4 award labels found: {found_awards} (0.25 pts)")
            total_score += 0.25
        elif match_count > 0:
            partial = round(match_count * points_per_label, 4)
            print(f"PARTIAL: Component 3 — {match_count}/4 labels found: {found_awards} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No award labels found. Text boxes: {label_texts}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Stars arranged in a row (similar top) and evenly spaced (0.15 points)
    try:
        if len(star_shapes) >= 4:
            # Check same row: all stars should have similar top values
            tops = [s.top for s in star_shapes]
            top_range = max(tops) - min(tops)
            # Allow tolerance of 0.5 inch = 457200 EMU
            same_row = top_range < 457200

            # Check even spacing: sort by left position, check gaps are similar
            sorted_stars = sorted(star_shapes, key=lambda s: s.left)
            lefts = [s.left for s in sorted_stars]
            gaps = [lefts[i+1] - lefts[i] for i in range(len(lefts)-1)]

            if len(gaps) >= 2:
                avg_gap = sum(gaps) / len(gaps)
                # Allow 20% tolerance on gap uniformity
                evenly_spaced = all(abs(g - avg_gap) / avg_gap < 0.2 for g in gaps) if avg_gap > 0 else False
            else:
                evenly_spaced = False

            if same_row and evenly_spaced:
                print(f"PASS: Component 4 — Stars in a row (top range={top_range}EMU) and evenly spaced (gaps={gaps}) (0.15 pts)")
                total_score += 0.15
            elif same_row:
                print(f"PARTIAL: Component 4 — Stars in a row but not evenly spaced (gaps={gaps}) (0.08 pts)")
                total_score += 0.08
            elif evenly_spaced:
                print(f"PARTIAL: Component 4 — Stars evenly spaced but not same row (tops={tops}) (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 4 — Stars not in a row (top range={top_range}) and not evenly spaced (gaps={gaps})")
        else:
            print(f"FAIL: Component 4 — Not enough stars ({len(star_shapes)}) to check layout")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
