"""
Initial Setup: Create a 10-slide sales pitch presentation
Task ID: impress_sales_084
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_084'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.paragraphs[0].text = point
        else:
            p = tf.add_paragraph()
            p.text = point
    return slide


def add_blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only layout
    return slide


def create_initial():
    prs = Presentation()

    # Slide 1: Title Slide
    add_title_slide(prs, "Accelerate Your Revenue Growth",
                    "CloudSync Solutions - Enterprise Sales Platform")

    # Slide 2: Company Overview
    add_content_slide(prs, "About CloudSync Solutions", [
        "Founded in 2019 with a mission to transform B2B sales",
        "Headquartered in San Francisco, offices in London and Singapore",
        "450+ employees across 12 countries",
        "Serving 2,800+ enterprise customers worldwide",
    ])

    # Slide 3: Market Opportunity
    add_content_slide(prs, "Market Opportunity", [
        "Global CRM market expected to reach $145.8B by 2028",
        "42% of sales teams report inefficient pipeline management",
        "Average sales rep spends 65% of time on non-selling activities",
        "AI-driven sales tools growing at 28% CAGR",
    ])

    # Slide 4: Product Suite
    add_content_slide(prs, "Product Suite Overview", [
        "CloudSync CRM - Intelligent customer relationship management",
        "Pipeline Pro - AI-powered deal forecasting and tracking",
        "Revenue Intelligence - Real-time analytics dashboard",
        "Engage360 - Multi-channel outreach automation",
    ])

    # Slide 5: Key Metrics
    add_content_slide(prs, "Key Performance Metrics", [
        "$78M ARR as of Q4 2024 (up 52% YoY)",
        "Net Revenue Retention Rate: 135%",
        "Average contract value: $42,000",
        "Customer acquisition cost payback: 11 months",
    ])

    # Slide 6: Customer Success Stories
    add_content_slide(prs, "Customer Success Stories", [
        "TechVault Inc. - Reduced sales cycle by 34% in 6 months",
        "Meridian Healthcare - 2.8x increase in qualified pipeline",
        "Nova Financial Group - 89% forecast accuracy improvement",
        "Atlas Manufacturing - $12M additional revenue in Year 1",
    ])

    # Slide 7: Competitive Advantages
    add_content_slide(prs, "Why CloudSync Wins", [
        "Proprietary AI engine trained on 50M+ sales interactions",
        "90-second onboarding with automatic CRM data migration",
        "Native integrations with 200+ enterprise tools",
        "SOC 2 Type II and ISO 27001 certified",
    ])

    # Slide 8: Growth Strategy
    add_content_slide(prs, "2025 Growth Strategy", [
        "Expand into APAC market with localized offerings",
        "Launch CloudSync Marketplace for partner integrations",
        "Release AI Sales Coach for real-time call guidance",
        "Target 150% ARR growth through enterprise upmarket push",
    ])

    # Slide 9: Blank placeholder (to be filled by agent with awards)
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout - NO title, NO content

    # Slide 10: Call to Action
    add_content_slide(prs, "Next Steps", [
        "Schedule a personalized demo with our solutions team",
        "Start a 30-day enterprise pilot program",
        "Contact: sales@cloudsync.io | +1 (415) 555-0198",
        "Visit cloudsync.io/enterprise for case studies",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
