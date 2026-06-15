"""
Reward Script: Apply unique background colors and speaker notes to each of 4 slides
Task ID: osworld_impress_note_bg_combined_009
Domain: libreoffice_impress

Scoring (8 sub-components, 0.125 each, total 1.0):
  - Slide 1 background is pale blue (#ADD8E6): 0.125
  - Slide 1 note is 'Welcome to the company': 0.125
  - Slide 2 background is pale yellow (#FFFFE0): 0.125
  - Slide 2 note is 'Overview of your first 30 days': 0.125
  - Slide 3 background is pale green (#90EE90): 0.125
  - Slide 3 note is 'Tools and systems walkthrough': 0.125
  - Slide 4 background is pale orange (#FFD580): 0.125
  - Slide 4 note is 'Meet your team and contacts': 0.125
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_note_bg_combined_009'

# Expected ground truth values from task_config context
EXPECTED = [
    {
        "slide": 1,
        "bg_hex": "ADD8E6",   # pale blue
        "note": "Welcome to the company",
    },
    {
        "slide": 2,
        "bg_hex": "FFFFE0",   # pale yellow
        "note": "Overview of your first 30 days",
    },
    {
        "slide": 3,
        "bg_hex": "90EE90",   # pale green
        "note": "Tools and systems walkthrough",
    },
    {
        "slide": 4,
        "bg_hex": "FFD580",   # pale orange
        "note": "Meet your team and contacts",
    },
]


def persist_app_state():
    """Send Ctrl+S to ensure any unsaved GUI edits are flushed to disk."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def get_slide_bg_hex(slide):
    """Return the hex string of the slide's solid background color, or None."""
    fill = slide.background.fill
    if fill.type == 1:  # SOLID
        try:
            return str(fill.fore_color.rgb).upper()
        except Exception:
            return None
    elif fill.type == 5:  # INHERITED from master
        try:
            master_fill = slide.slide_layout.slide_master.background.fill
            if master_fill.type == 1:
                return str(master_fill.fore_color.rgb).upper()
        except Exception:
            pass
    return None


def get_slide_notes_text(slide):
    """Return stripped notes text from the slide, or empty string on error."""
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition gate: must have exactly 4 slides
    num_slides = len(prs.slides)
    if num_slides != 4:
        print(f"CRITICAL: Expected 4 slides, found {num_slides}. Cannot score.")
        print("REWARD: 0.0")
        return 0.0

    for spec in EXPECTED:
        slide_idx = spec["slide"] - 1  # 0-based index
        slide = prs.slides[slide_idx]
        slide_num = spec["slide"]

        # Sub-component A: Background color (0.125 points)
        try:
            actual_bg = get_slide_bg_hex(slide)
            expected_bg = spec["bg_hex"].upper()
            if actual_bg == expected_bg:
                print(f"PASS: Slide {slide_num} background = #{actual_bg} (expected #{expected_bg}) (0.125 pts)")
                total_score += 0.125
            else:
                print(f"FAIL: Slide {slide_num} background = #{actual_bg}, expected #{expected_bg}")
        except Exception as e:
            print(f"ERROR: Slide {slide_num} background check failed: {e}")

        # Sub-component B: Speaker note (0.125 points)
        try:
            actual_note = get_slide_notes_text(slide)
            expected_note = spec["note"]
            if actual_note == expected_note:
                print(f"PASS: Slide {slide_num} note = {repr(actual_note)} (0.125 pts)")
                total_score += 0.125
            else:
                print(f"FAIL: Slide {slide_num} note = {repr(actual_note)}, expected {repr(expected_note)}")
        except Exception as e:
            print(f"ERROR: Slide {slide_num} note check failed: {e}")

    final_score = round(min(total_score, 1.0), 4)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
