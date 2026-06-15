"""
Initial Setup: Employee onboarding training deck (4 slides, white backgrounds, no notes)
Task ID: osworld_impress_note_bg_combined_009
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_note_bg_combined_009'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard widescreen dimensions (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ----------------------------------------------------------------
    # Slide 1: Welcome / Introduction
    # ----------------------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout

    # White background — explicit solid white fill
    bg1 = slide1.background.fill
    bg1.solid()
    bg1.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    slide1.shapes.title.text = "Welcome to AcmeCorp"
    slide1.placeholders[1].text = "New Hire Onboarding Program — 2025 Cohort"

    # ----------------------------------------------------------------
    # Slide 2: First 30 Days Overview
    # ----------------------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content

    bg2 = slide2.background.fill
    bg2.solid()
    bg2.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    slide2.shapes.title.text = "Your First 30 Days"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Week 1: Orientation & HR Paperwork"
    p2b = tf2.add_paragraph()
    p2b.text = "Week 2: Department Introductions & Shadow Sessions"
    p2c = tf2.add_paragraph()
    p2c.text = "Week 3: Hands-On Project Assignments"
    p2d = tf2.add_paragraph()
    p2d.text = "Week 4: First Performance Check-In with Manager"

    # ----------------------------------------------------------------
    # Slide 3: Tools and Systems
    # ----------------------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content

    bg3 = slide3.background.fill
    bg3.solid()
    bg3.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    slide3.shapes.title.text = "Tools & Systems Walkthrough"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Slack — Team communication platform"
    p3b = tf3.add_paragraph()
    p3b.text = "Jira — Project tracking and sprint planning"
    p3c = tf3.add_paragraph()
    p3c.text = "Confluence — Internal documentation wiki"
    p3d = tf3.add_paragraph()
    p3d.text = "Workday — HR, payroll, and benefits portal"
    p3e = tf3.add_paragraph()
    p3e.text = "GitHub Enterprise — Version control and code reviews"

    # ----------------------------------------------------------------
    # Slide 4: Meet the Team
    # ----------------------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content

    bg4 = slide4.background.fill
    bg4.solid()
    bg4.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    slide4.shapes.title.text = "Meet Your Team & Contacts"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Sarah Chen — Engineering Lead (sarah.chen@acmecorp.com)"
    p4b = tf4.add_paragraph()
    p4b.text = "Marcus Johnson — People Operations (marcus.j@acmecorp.com)"
    p4c = tf4.add_paragraph()
    p4c.text = "Priya Sharma — IT Helpdesk (it-support@acmecorp.com)"
    p4d = tf4.add_paragraph()
    p4d.text = "David Park — Facilities & Security (facilities@acmecorp.com)"

    # Save — NO speaker notes on any slide
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
