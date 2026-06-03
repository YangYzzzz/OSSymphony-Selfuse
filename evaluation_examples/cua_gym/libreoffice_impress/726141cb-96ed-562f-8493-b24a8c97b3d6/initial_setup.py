"""
Initial Setup: Create Teaching_Plan presentation with 8 slides, no notes on slides 2-6.
Task ID: impress_teach_026
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_026'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=None, alignment=None):
    """Helper to add a text box to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullet_points(slide, left, top, width, height, items, font_size=16):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(font_size)
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    add_text_box(slide1, 1.5, 1.5, 10, 1.5, "Teaching Plan: Introduction to Data Science",
                 font_size=36, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, 2, 3.5, 9, 1, "Instructor: Dr. Amelia Richardson",
                 font_size=22, color=RGBColor(0xCC, 0xDD, 0xEE), alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, 2, 4.5, 9, 1, "Spring Semester 2025 | Module 3: Statistical Foundations",
                 font_size=18, color=RGBColor(0xAA, 0xBB, 0xCC), alignment=PP_ALIGN.CENTER)

    # --- Slide 2: Learning Objectives ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide2, 0.8, 0.4, 10, 0.8, "Learning Objectives",
                 font_size=32, bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))
    add_bullet_points(slide2, 0.8, 1.5, 11, 5, [
        "Understand the fundamentals of descriptive statistics",
        "Apply measures of central tendency to real-world datasets",
        "Interpret standard deviation and variance in context",
        "Distinguish between population and sample statistics",
        "Construct and read basic histograms and box plots",
    ], font_size=18)
    # NO notes on slide 2

    # --- Slide 3: Key Concepts Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide3, 0.8, 0.4, 10, 0.8, "Key Concepts Overview",
                 font_size=32, bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))
    add_bullet_points(slide3, 0.8, 1.5, 5.5, 5, [
        "Mean, Median, Mode",
        "Range and Interquartile Range (IQR)",
        "Variance and Standard Deviation",
        "Normal Distribution Properties",
        "Skewness and Kurtosis",
    ], font_size=18)
    add_bullet_points(slide3, 7, 1.5, 5.5, 5, [
        "Sampling techniques overview",
        "Confidence intervals basics",
        "Hypothesis testing introduction",
        "P-value interpretation",
        "Type I and Type II errors",
    ], font_size=18)
    # NO notes on slide 3

    # --- Slide 4: Classroom Activity ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide4, 0.8, 0.4, 10, 0.8, "Classroom Activity: Analyzing Survey Data",
                 font_size=32, bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))
    add_bullet_points(slide4, 0.8, 1.5, 11, 4.5, [
        "Activity Duration: 20 minutes",
        "Students will work with the Campus Life Survey dataset (N=450)",
        "Calculate mean, median, and mode for student commute times",
        "Compare results across three campus locations",
        "Identify potential outliers using the 1.5 IQR rule",
        "Present findings to a partner and discuss discrepancies",
    ], font_size=18)
    # NO notes on slide 4

    # --- Slide 5: Statistical Diagrams ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide5, 0.8, 0.4, 10, 0.8, "Statistical Diagrams and Visualization",
                 font_size=32, bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))
    add_bullet_points(slide5, 0.8, 1.5, 5.5, 5, [
        "Box Plot Anatomy:",
        "  - Whiskers show data range",
        "  - Box spans Q1 to Q3",
        "  - Line inside box = median",
        "  - Dots beyond whiskers = outliers",
    ], font_size=16)
    add_bullet_points(slide5, 7, 1.5, 5.5, 5, [
        "Histogram Guidelines:",
        "  - Choose appropriate bin width",
        "  - Label axes clearly",
        "  - Show frequency or relative frequency",
        "  - Note distribution shape",
    ], font_size=16)
    # NO notes on slide 5

    # --- Slide 6: Summary and Q&A ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide6, 0.8, 0.4, 10, 0.8, "Summary and Review",
                 font_size=32, bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))
    add_bullet_points(slide6, 0.8, 1.5, 11, 4.5, [
        "Today we covered: descriptive statistics fundamentals",
        "Key formulas: mean, variance, standard deviation",
        "Practical application through the Campus Life Survey analysis",
        "Next session: Inferential statistics and hypothesis testing",
        "Reading assignment: Chapter 5 (Sections 5.1 - 5.4)",
        "Problem set due: Friday, March 28, 2025",
    ], font_size=18)
    # NO notes on slide 6

    # --- Slide 7: Homework Assignments ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide7, 0.8, 0.4, 10, 0.8, "Homework and Assignments",
                 font_size=32, bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))
    add_bullet_points(slide7, 0.8, 1.5, 11, 4.5, [
        "Problem Set 3: Due March 28, 2025",
        "  - Questions 1-5: Calculate descriptive statistics from provided dataset",
        "  - Questions 6-8: Create box plots using spreadsheet software",
        "  - Question 9: Written analysis (250 words minimum)",
        "Group Project Milestone 2: Due April 4, 2025",
        "  - Submit data collection plan and preliminary analysis",
    ], font_size=18)

    # --- Slide 8: Resources ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide8, 0.8, 0.4, 10, 0.8, "Additional Resources",
                 font_size=32, bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))
    add_bullet_points(slide8, 0.8, 1.5, 11, 4.5, [
        "Textbook: 'Statistics for Data Science' by Thompson & Park, 4th Edition",
        "Online: Khan Academy - Statistics and Probability Module",
        "Software: R Studio (free download) or Python with pandas/matplotlib",
        "Office Hours: Tuesdays 2-4 PM, Room 312 Science Building",
        "Teaching Assistant: James Okoro (james.okoro@university.edu)",
        "Discussion Forum: Canvas > Module 3 Discussion Board",
    ], font_size=18)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
