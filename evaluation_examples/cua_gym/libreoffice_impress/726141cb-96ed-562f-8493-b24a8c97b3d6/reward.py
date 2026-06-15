"""
Reward Script: Add presenter notes to slides 2-6 with specific text
Task ID: impress_teach_026
Domain: libreoffice_impress
Scoring: 5 components (0.2 each) - one per slide note verification
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_026'

# Expected notes for slides 2-6 (0-indexed: slides[1] through slides[5])
EXPECTED_NOTES = {
    2: "Spend 5 minutes on introductory concepts.",
    3: "Show video clip at this point.",
    4: "Ask students to discuss in pairs.",
    5: "Review the diagram carefully.",
    6: "Summarize and take questions.",
}

POINTS_PER_SLIDE = 0.2


def persist_app_state():
    """Try to save any unsaved LibreOffice edits before verification."""
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    for slide_num, expected_text in EXPECTED_NOTES.items():
        slide_idx = slide_num - 1  # 0-based index
        try:
            slide = prs.slides[slide_idx]
            # Get notes text
            try:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                notes_text = ""

            if notes_text == expected_text:
                print(f"PASS: Slide {slide_num} notes match exactly ({POINTS_PER_SLIDE} pts)")
                total_score += POINTS_PER_SLIDE
            elif expected_text.lower() in notes_text.lower():
                # Partial: the expected text is contained but not exact
                partial = POINTS_PER_SLIDE * 0.5
                print(f"PARTIAL: Slide {slide_num} notes contain expected text but not exact match ({partial} pts)")
                print(f"  Expected: {expected_text!r}")
                print(f"  Found:    {notes_text!r}")
                total_score += partial
            else:
                print(f"FAIL: Slide {slide_num} notes do not match")
                print(f"  Expected: {expected_text!r}")
                print(f"  Found:    {notes_text!r}")
        except Exception as e:
            print(f"ERROR: Slide {slide_num} check failed: {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist unsaved GUI state before verification
persist_app_state()

# Run verification
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
