"""
Reward Script: Replace last slide with departmental closing slide
Task ID: impress_cross_acad_037
Domain: libreoffice_impress
Scoring:
  Component 1: Last slide has white background (0.10 pts)
  Component 2: 'Department of Bioinformatics' text present on last slide (0.10 pts)
  Component 2b: 'Department of Bioinformatics' is 36pt and bold (0.15 pts)
  Component 2c: 'Department of Bioinformatics' is color #003087 and centered (0.10 pts)
  Component 3: 'University of Cambridge' text present on last slide (0.10 pts)
  Component 3b: 'University of Cambridge' is 24pt and centered (0.15 pts)
  Component 4a: Image (logo) is present on last slide (0.10 pts)
  Component 4b: Image is horizontally centered (0.10 pts)
  Component 4c: Image vertical position is near y=14cm (0.10 pts)
  Total: 1.00
"""

import os
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_cross_acad_037'


def get_all_text_in_slide(slide):
    """Return all text strings visible in a slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return texts


def find_paragraph_with_text(slide, search_text):
    """Find a (shape, paragraph) pair containing the given text."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if search_text in para.text:
                    return (shape, para)
    return (None, None)


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file must have exactly 20 slides
    num_slides = len(prs.slides)
    if num_slides != 20:
        print(f"FAIL: Expected 20 slides, found {num_slides}. Task was to replace last slide.")
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    last_slide = prs.slides[19]  # slide index 19 = slide 20

    # Precondition: old 'Thank You' content must be gone (gate — no points)
    all_text = get_all_text_in_slide(last_slide)
    old_content_keywords = ['Thank You', 'Questions and Discussion', 'elena.marchetti', 'marchetti-lab']
    old_content_present = any(
        any(kw.lower() in t.lower() for kw in old_content_keywords)
        for t in all_text
    )
    if old_content_present:
        print(f"FAIL (gate): Old 'Thank You' slide content still present on slide 20")
        print(f"  Found: {all_text[:3]}")
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 1: White background on last slide (0.10 points)
    # Task specifies "white background"; golden file uses SOLID fill FFFFFF
    try:
        fill = last_slide.background.fill
        # Accept explicit white (FFFFFF) or inherited/None (renders as white)
        if fill.type == 1 and str(fill.fore_color.rgb).upper() == 'FFFFFF':
            print(f"PASS: Component 1 — White background (FFFFFF) confirmed (0.10 pts)")
            total_score += 0.10
        elif fill.type is None:
            print(f"PASS: Component 1 — Background is inherited/default (appears white) (0.10 pts)")
            total_score += 0.10
        else:
            fill_color = str(fill.fore_color.rgb).upper() if fill.type == 1 else 'unknown'
            print(f"FAIL: Component 1 — Expected white background, fill type={fill.type}, color={fill_color}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: 'Department of Bioinformatics' text present (0.10 points)
    dept_para = None
    dept_runs = []
    try:
        _, dept_para = find_paragraph_with_text(last_slide, 'Department of Bioinformatics')
        if dept_para is not None:
            print(f"PASS: Component 2 — 'Department of Bioinformatics' text found (0.10 pts)")
            total_score += 0.10
            dept_runs = [r for r in dept_para.runs if (r.text or '').strip()]
        else:
            print(f"FAIL: Component 2 — 'Department of Bioinformatics' not found on last slide")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 2b: 'Department of Bioinformatics' — 36pt bold (0.15 points)
    try:
        if dept_para is not None and dept_runs:
            run = dept_runs[0]
            font_size_pt = run.font.size / 12700 if run.font.size else None
            is_bold = run.font.bold is True
            size_ok = font_size_pt is not None and abs(font_size_pt - 36.0) < 0.5
            if size_ok and is_bold:
                print(f"PASS: Component 2b — 36pt bold confirmed (size={font_size_pt:.1f}pt, bold=True) (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 2b — Expected 36pt bold; found size={font_size_pt}, bold={run.font.bold}")
        elif dept_para is not None:
            print(f"FAIL: Component 2b — No non-empty runs in 'Department of Bioinformatics' paragraph")
    except Exception as e:
        print(f"ERROR: Component 2b — {e}")

    # Component 2c: 'Department of Bioinformatics' — color #003087 and CENTER (0.10 points)
    try:
        if dept_para is not None and dept_runs:
            run = dept_runs[0]
            color_ok = False
            try:
                if run.font.color.type is not None:
                    rgb_str = str(run.font.color.rgb).upper()
                    color_ok = (rgb_str == '003087')
                    if not color_ok:
                        print(f"FAIL: Component 2c — Expected color #003087, found #{rgb_str}")
                else:
                    print(f"FAIL: Component 2c — Font color is inherited (type=None), expected #003087")
            except Exception as ce:
                print(f"FAIL: Component 2c — Could not read color: {ce}")

            align_ok = (dept_para.alignment == PP_ALIGN.CENTER)
            if not align_ok:
                print(f"FAIL: Component 2c — Expected CENTER alignment, found {dept_para.alignment}")

            if color_ok and align_ok:
                print(f"PASS: Component 2c — Color #003087 and CENTER alignment confirmed (0.10 pts)")
                total_score += 0.10
        elif dept_para is not None:
            print(f"FAIL: Component 2c — No non-empty runs to check color/alignment")
    except Exception as e:
        print(f"ERROR: Component 2c — {e}")

    # Component 3: 'University of Cambridge' text present (0.10 points)
    univ_para = None
    univ_runs = []
    try:
        _, univ_para = find_paragraph_with_text(last_slide, 'University of Cambridge')
        if univ_para is not None:
            print(f"PASS: Component 3 — 'University of Cambridge' text found (0.10 pts)")
            total_score += 0.10
            univ_runs = [r for r in univ_para.runs if (r.text or '').strip()]
        else:
            print(f"FAIL: Component 3 — 'University of Cambridge' not found on last slide")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 3b: 'University of Cambridge' — 24pt and CENTER (0.15 points)
    try:
        if univ_para is not None and univ_runs:
            run = univ_runs[0]
            font_size_pt = run.font.size / 12700 if run.font.size else None
            size_ok = font_size_pt is not None and abs(font_size_pt - 24.0) < 0.5
            align_ok = (univ_para.alignment == PP_ALIGN.CENTER)

            if size_ok and align_ok:
                print(f"PASS: Component 3b — 24pt CENTER confirmed (size={font_size_pt:.1f}pt) (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 3b — Expected 24pt CENTER; found size={font_size_pt}, align={univ_para.alignment}")
        elif univ_para is not None:
            print(f"FAIL: Component 3b — No non-empty runs in 'University of Cambridge' paragraph")
    except Exception as e:
        print(f"ERROR: Component 3b — {e}")

    # Component 4a: Image (logo) present on last slide (0.10 points)
    logo_shape = None
    try:
        for shape in last_slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                logo_shape = shape
                break
        if logo_shape is not None:
            print(f"PASS: Component 4a — Image found on last slide (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 4a — No image found on last slide")
    except Exception as e:
        print(f"ERROR: Component 4a — {e}")

    # Component 4b: Image horizontally centered (0.10 points)
    try:
        if logo_shape is not None:
            slide_cx = prs.slide_width / 2
            shape_cx = logo_shape.left + logo_shape.width / 2
            # Allow 1% relative tolerance for centering
            horiz_ok = abs(shape_cx - slide_cx) / max(abs(slide_cx), 1) <= 0.01
            if horiz_ok:
                print(f"PASS: Component 4b — Image centered horizontally "
                      f"(shape_cx={shape_cx/914400:.3f}in, slide_cx={slide_cx/914400:.3f}in) (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 4b — Image not centered "
                      f"(shape_cx={shape_cx/914400:.3f}in, slide_cx={slide_cx/914400:.3f}in)")
    except Exception as e:
        print(f"ERROR: Component 4b — {e}")

    # Component 4c: Image vertical position near y=14cm (0.10 points)
    # 14 cm = 14 * 360000 = 5040000 EMU
    try:
        if logo_shape is not None:
            target_y_emu = 14 * 360000  # 5040000 EMU
            actual_y_emu = logo_shape.top
            # Allow ±5% tolerance (~0.7cm)
            y_tolerance = target_y_emu * 0.05
            y_ok = abs(actual_y_emu - target_y_emu) <= y_tolerance
            if y_ok:
                print(f"PASS: Component 4c — Image at y={actual_y_emu/360000:.2f}cm (target 14cm, tol ±5%) (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 4c — Image at y={actual_y_emu/360000:.2f}cm, expected ~14cm (±5%)")
    except Exception as e:
        print(f"ERROR: Component 4c — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against golden file (pptx version on VM)
file_path = f'{WORKDIR}/{TASK_ID}_initial.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
