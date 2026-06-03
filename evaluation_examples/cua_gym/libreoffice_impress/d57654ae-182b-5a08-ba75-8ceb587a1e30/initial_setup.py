"""
Initial Setup: Bioinformatics seminar presentation with 20 slides
Task ID: impress_cross_acad_037
Domain: libreoffice_impress
Creates: ~/Documents/bioinformatics_seminar.odp (20 slides, last has custom Thank You content)
         ~/Pictures/dept_logo.png (departmental logo)
"""

import os
import subprocess
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont
import io

WORKDIR = '/home/user'
TASK_ID = 'impress_cross_acad_037'
DOCS_DIR = f'{WORKDIR}/Documents'
PICS_DIR = f'{WORKDIR}/Pictures'
PPTX_PATH = f'{WORKDIR}/{TASK_ID}_initial.pptx'
ODP_PATH = f'{DOCS_DIR}/bioinformatics_seminar.odp'


def create_logo_image():
    """Create a realistic departmental logo PNG."""
    os.makedirs(PICS_DIR, exist_ok=True)

    # Create a 400x200 logo image with blue background and text
    img = Image.new('RGBA', (400, 200), (255, 255, 255, 255))
    draw = ImageDraw.Draw(img)

    # Draw a blue rectangle as the logo background
    draw.rectangle([20, 20, 380, 180], fill=(0, 48, 135, 255), outline=(0, 0, 0, 255), width=2)

    # Draw a circular crest shape
    draw.ellipse([160, 30, 240, 110], fill=(255, 255, 255, 255), outline=(0, 48, 135, 255), width=2)

    # Add text elements (using default font since custom may not be available)
    # "Dept of Bioinformatics" text area
    draw.rectangle([30, 115, 370, 145], fill=(255, 255, 255, 100))

    # Draw DNA helix-like lines
    for i in range(10):
        x1 = 170 + int(20 * ((-1) ** i) * 0.5)
        y1 = 35 + i * 7
        x2 = 230 - int(20 * ((-1) ** i) * 0.5)
        y2 = y1
        draw.line([(x1, y1), (x2, y2)], fill=(0, 100, 200, 255), width=2)
        if i < 9:
            draw.line([(200, y1), (200, y1 + 7)], fill=(0, 200, 100, 255), width=1)

    logo_path = f'{PICS_DIR}/dept_logo.png'
    img.save(logo_path, 'PNG')
    print(f'Logo created: {logo_path}')
    return logo_path


def create_presentation():
    """Create a 20-slide bioinformatics seminar presentation."""
    os.makedirs(DOCS_DIR, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide layout indices in default template
    TITLE_SLIDE_LAYOUT = 0    # Title slide
    TITLE_CONTENT_LAYOUT = 1  # Title + Content
    BLANK_LAYOUT = 6          # Blank (may vary)

    # ---- Slide 1: Title Slide ----
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Advances in Comparative Genomics"
    subtitle = slide.placeholders[1]
    subtitle.text = "Department of Bioinformatics Seminar Series\nDr. Elena Marchetti\nUniversity of Cambridge\nMarch 2025"

    # ---- Slide 2: Outline ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Seminar Outline"
    tf = slide.placeholders[1].text_frame
    tf.text = "1. Introduction to Comparative Genomics"
    items = [
        "2. Sequence Alignment Methods",
        "3. Phylogenetic Analysis Tools",
        "4. Gene Family Evolution",
        "5. Case Study: Primate Genomes",
        "6. Machine Learning Applications",
        "7. Future Directions",
        "8. Q&A Session",
    ]
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    # ---- Slide 3: Introduction ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Introduction to Comparative Genomics"
    tf = slide.placeholders[1].text_frame
    tf.text = "What is Comparative Genomics?"
    for bullet in [
        "Study of genomic features by comparing genomes of different species",
        "Identifies conserved regions, gene families, and evolutionary changes",
        "Fundamental to understanding disease mechanisms and drug targets",
        "Enabled by decreasing sequencing costs and improved algorithms",
    ]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1

    # ---- Slide 4: Historical Context ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Historical Context"
    tf = slide.placeholders[1].text_frame
    tf.text = "Milestones in Comparative Genomics"
    for bullet in [
        "1995: First bacterial genome sequenced (H. influenzae)",
        "2001: Human Genome Project draft complete",
        "2003: Mouse genome comparison revealed conserved elements",
        "2012: ENCODE project identifies functional DNA elements",
        "2020: Pangenome reference assemblies introduced",
    ]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1

    # ---- Slide 5: Sequence Alignment Methods ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Sequence Alignment Methods"
    tf = slide.placeholders[1].text_frame
    tf.text = "Core Algorithms"
    for bullet in [
        "Pairwise alignment: Smith-Waterman, Needleman-Wunsch",
        "Multiple sequence alignment: ClustalW, MUSCLE, MAFFT",
        "Whole genome alignment: LASTZ, MAUVE, Mugsy",
        "Key parameters: gap penalties, substitution matrices (BLOSUM, PAM)",
        "Trade-off between sensitivity and computational efficiency",
    ]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1

    # ---- Slide 6: Alignment Scoring ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Alignment Scoring Matrices"
    tf = slide.placeholders[1].text_frame
    tf.text = "BLOSUM and PAM Matrices"
    for bullet in [
        "BLOSUM62: derived from conserved domain alignments",
        "PAM250: based on point accepted mutations model",
        "Higher BLOSUM numbers = more divergent sequences",
        "Matrix selection impacts alignment sensitivity significantly",
    ]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1

    # ---- Slide 7: Phylogenetic Analysis ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Phylogenetic Analysis Tools"
    tf = slide.placeholders[1].text_frame
    tf.text = "Methods and Software"
    for bullet in [
        "Distance-based: UPGMA, Neighbor-Joining",
        "Maximum Parsimony: finds tree with fewest evolutionary changes",
        "Maximum Likelihood: RAxML, IQ-TREE",
        "Bayesian inference: MrBayes, BEAST",
        "Bootstrap resampling for statistical support values",
    ]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1

    # ---- Slide 8: Tree Interpretation ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Interpreting Phylogenetic Trees"
    tf = slide.placeholders[1].text_frame
    tf.text = "Key Concepts"
    for bullet in [
        "Nodes represent common ancestors",
        "Branch lengths indicate evolutionary distance",
        "Bootstrap values >= 70% indicate strong support",
        "Outgroup selection affects tree topology",
        "Molecular clock assumption: uniform mutation rate",
    ]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1

    # ---- Slide 9: Gene Family Evolution ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Gene Family Evolution"
    tf = slide.placeholders[1].text_frame
    tf.text = "Mechanisms of Gene Family Expansion"
    for bullet in [
        "Gene duplication: tandem, segmental, whole-genome",
        "Divergence: subfunctionalization and neofunctionalization",
        "Gene loss: pseudogenization and deletion",
        "Horizontal gene transfer in prokaryotes",
        "Tools: OrthoFinder, OrthoMCL, CAFE",
    ]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1

    # ---- Slide 10: Case Study - Primates ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Case Study: Primate Genome Comparison"
    tf = slide.placeholders[1].text_frame
    tf.text = "Human-Chimpanzee Divergence"
    for bullet in [
        "~98.7% DNA sequence identity with chimpanzees",
        "~45 million SNPs between human and chimp genomes",
        "Accelerated Human Regions (AHRs): 202 identified regions",
        "Key differences in gene regulation, not just coding sequence",
        "FOXP2: language-associated gene shows rapid evolution",
    ]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1

    # ---- Slide 11: Structural Variation ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Structural Variation in Primate Genomes"
    tf = slide.placeholders[1].text_frame
    tf.text = "Types of Structural Variants"
    for bullet in [
        "Copy Number Variants (CNVs): >1kb duplications/deletions",
        "Inversions: segments flipped in orientation",
        "Translocations: segments moved between chromosomes",
        "Insertions of mobile elements (SINE, LINE, LTR)",
        "AMY1 gene: human copy number expansion for starch digestion",
    ]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1

    # ---- Slide 12: Machine Learning in Genomics ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Machine Learning Applications"
    tf = slide.placeholders[1].text_frame
    tf.text = "Deep Learning for Sequence Analysis"
    for bullet in [
        "CNN-based models: DeepBind, DeepSEA for regulatory elements",
        "Transformer models: DNABERT, Nucleotide Transformer",
        "Graph Neural Networks for protein interaction networks",
        "Predicting gene expression from sequence alone",
        "Cross-species transfer learning improves rare species analysis",
    ]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1

    # ---- Slide 13: Regulatory Element Conservation ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Regulatory Element Conservation"
    tf = slide.placeholders[1].text_frame
    tf.text = "Non-Coding Conserved Elements"
    for bullet in [
        "Ultra-conserved elements (UCEs): identical in human/mouse/rat",
        "Enhancers show species-specific activity patterns",
        "VISTA browser: visualizing vertebrate conservation",
        "~5% of human genome under purifying selection",
        "Many conserved elements have unknown function",
    ]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1

    # ---- Slide 14: Proteomics Integration ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Integrating Proteomics Data"
    tf = slide.placeholders[1].text_frame
    tf.text = "Multi-Omics Comparative Analysis"
    for bullet in [
        "Protein conservation often higher than DNA sequence",
        "Mass spectrometry identifies conserved protein complexes",
        "Interactome comparisons reveal network rewiring",
        "AlphaFold2: structural predictions across species",
        "Structural homologs may lack sequence similarity",
    ]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1

    # ---- Slide 15: Epigenomics ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Comparative Epigenomics"
    tf = slide.placeholders[1].text_frame
    tf.text = "Epigenetic Regulation Across Species"
    for bullet in [
        "DNA methylation patterns show evolutionary conservation",
        "Histone modifications partially conserved in mammals",
        "ATAC-seq reveals species-specific open chromatin regions",
        "Imprinting: conserved in placental mammals",
        "Epigenetic clock: universal methylation-based age predictor",
    ]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1

    # ---- Slide 16: Tools and Databases ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Databases and Tools"
    tf = slide.placeholders[1].text_frame
    tf.text = "Essential Resources"
    for bullet in [
        "UCSC Genome Browser: multi-species alignment tracks",
        "Ensembl: gene annotations and synteny analysis",
        "NCBI RefSeq: curated reference sequences",
        "STRING: protein interaction networks",
        "Galaxy: workflow platform for reproducible analysis",
    ]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1

    # ---- Slide 17: Current Challenges ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Current Challenges"
    tf = slide.placeholders[1].text_frame
    tf.text = "Open Problems in the Field"
    for bullet in [
        "Genome assemblies remain fragmented for many species",
        "Annotation of non-coding RNA across species",
        "Distinguishing adaptive from neutral evolution",
        "Computational scalability for pangenome analysis",
        "Integrating long-read sequencing data effectively",
    ]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1

    # ---- Slide 18: Future Directions ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Future Directions"
    tf = slide.placeholders[1].text_frame
    tf.text = "Emerging Opportunities"
    for bullet in [
        "Earth BioGenome Project: sequence all eukaryotic life",
        "Single-cell comparative transcriptomics",
        "Long-read assemblies resolving complex regions",
        "AI-driven functional annotation of genomes",
        "Comparative spatial genomics: tissue-specific conservation",
    ]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1

    # ---- Slide 19: Acknowledgements ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Acknowledgements"
    tf = slide.placeholders[1].text_frame
    tf.text = "This research was supported by:"
    for bullet in [
        "Wellcome Trust Grant #218328",
        "BBSRC Research Grant BB/T012234/1",
        "ERC Advanced Grant 789123-CompGenomics",
        "Collaborators: Dr. Priya Sharma (Oxford), Prof. Lars Eriksen (Copenhagen)",
        "Computational resources: Cambridge HPC Service",
    ]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1

    # ---- Slide 20: Last slide (custom Thank You content - this will be replaced by agent) ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Thank You"
    tf = slide.placeholders[1].text_frame
    tf.text = "Questions and Discussion"
    for bullet in [
        "Contact: elena.marchetti@bioinf.cam.ac.uk",
        "Lab website: http://marchetti-lab.bioinf.cam.ac.uk",
        "Twitter/X: @MarchettiLab",
        "Preprint: bioRxiv 2025.01.15.634321",
    ]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1

    # Add background color to last slide (custom style to make it distinct)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF0, 0xF8, 0xFF)  # alice blue background

    # Save as pptx first
    prs.save(PPTX_PATH)
    print(f'PPTX created: {PPTX_PATH}')

    # Convert pptx to odp using LibreOffice
    result = subprocess.run(
        ['libreoffice', '--headless', '--convert-to', 'odp', '--outdir', DOCS_DIR, PPTX_PATH],
        capture_output=True, text=True, timeout=60
    )
    if result.returncode == 0:
        # LibreOffice converts and saves as <TASK_ID>_initial.odp in DOCS_DIR
        converted = f'{DOCS_DIR}/{TASK_ID}_initial.odp'
        # Rename to the desired filename
        if os.path.exists(converted):
            os.rename(converted, ODP_PATH)
            print(f'Converted and renamed to: {ODP_PATH}')
        else:
            print(f'Conversion output: {result.stdout}')
            print(f'Conversion error: {result.stderr}')
    else:
        print(f'LibreOffice conversion failed: {result.stderr}')
        print(f'Trying alternative: copying pptx to odp path')
        import shutil
        shutil.copy(PPTX_PATH, ODP_PATH)
        print(f'Copied pptx as odp: {ODP_PATH}')

    return PPTX_PATH


def main():
    print('Creating departmental logo...')
    create_logo_image()

    print('Creating bioinformatics seminar presentation...')
    create_presentation()

    print('Verifying files...')
    for path in [ODP_PATH, f'{PICS_DIR}/dept_logo.png', PPTX_PATH]:
        if os.path.exists(path):
            size = os.path.getsize(path)
            print(f'  OK: {path} ({size} bytes)')
        else:
            print(f'  MISSING: {path}')


main()
