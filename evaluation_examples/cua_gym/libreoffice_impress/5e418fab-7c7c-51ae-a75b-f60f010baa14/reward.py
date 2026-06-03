"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 241 the heading disappears because the background is still the default light grey. I want the slide to pop: please change the slide’s background fill to the built-in color “Dark Red 2” (hex #7F0000) and then set the title text itself to pure white #FFFFFF so it stands out.
Generated: 2025-09-10 20:56:43
Status: success
Model: azure-o3
Total Steps: 7
"""

import os
from pptx import Presentation
from pptx.enum.dml import MSO_FILL, MSO_COLOR_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER


def _rgb_to_hex(rgb_obj):
    """Convert a python-pptx RGBColor into a 6-char upper-case hex string (e.g. FFFFFF)."""
    if rgb_obj is None:
        return None
    # python-pptx RGBColor has __str__ yielding the hex value
    hex_str = str(rgb_obj).upper()
    # Strip a leading 0X if ever present and keep last 6 chars
    if hex_str.startswith("0X"):
        hex_str = hex_str[2:]
    return hex_str[-6:]


def _verify_background_color(slide, target_hex: str) -> bool:
    """Return True if slide background is solid fill of target_hex RGB."""
    fill = slide.background.fill
    if fill.type != MSO_FILL.SOLID:
        print(f"✗ Background fill is not solid (type={fill.type})")
        return False
    rgb_hex = _rgb_to_hex(fill.fore_color.rgb)
    print(f"Background RGB detected: {rgb_hex}")
    if rgb_hex == target_hex.upper():
        print("✓ Background color matches target")
        return True
    print("✗ Background color does not match target")
    return False


def _verify_title_color(slide, target_hex: str) -> bool:
    """Return True if ≥80 % of explicitly coloured title runs match target_hex."""
    title_phs = [sh for sh in slide.shapes
                 if sh.is_placeholder and sh.placeholder_format.type in (PP_PLACEHOLDER.TITLE,
                                                                          PP_PLACEHOLDER.CENTER_TITLE)]
    if not title_phs:
        print("✗ No title placeholder found on the slide")
        return False

    explicit_runs = 0
    matching_runs = 0
    for ph in title_phs:
        if not ph.has_text_frame:
            continue
        for para in ph.text_frame.paragraphs:
            for run in para.runs:
                col = run.font.color
                if col.type == MSO_COLOR_TYPE.RGB:  # only consider explicitly set RGB runs
                    explicit_runs += 1
                    if _rgb_to_hex(col.rgb) == target_hex.upper():
                        matching_runs += 1
    if explicit_runs == 0:
        print("✗ No explicit RGB-coloured runs in title text – cannot verify")
        return False

    ratio = matching_runs / explicit_runs
    print(f"Title colour matches: {matching_runs}/{explicit_runs} runs (ratio {ratio:.2f})")
    if ratio >= 0.8:  # tolerance for small discrepancies
        print("✓ Title text colour matches target")
        return True
    print("✗ Title text colour does not sufficiently match target")
    return False


def verify_task(file_path: str) -> float:
    """Verify that slide 241 has Dark Red 2 background (#7F0000) and white title text (#FFFFFF).

    Returns a progressive score between 0.0 and 1.0.
    """
    print(f"Verifying presentation: {file_path}")

    # ---------- Preliminary checks ----------
    if not os.path.exists(file_path):
        print("✗ File does not exist")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Could not load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    target_index = 240  # zero-based index for slide 241
    if len(prs.slides) <= target_index:
        print(f"✗ Presentation only has {len(prs.slides)} slides; need ≥241")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[target_index]

    # ---------- Requirement verifications ----------
    score = 0.0

    # 1. Background colour (#7F0000)
    if _verify_background_color(slide, "7F0000"):
        score += 0.5

    # 2. Title text colour (#FFFFFF)
    if _verify_title_color(slide, "FFFFFF"):
        score += 0.5

    # ---------- Final score ----------
    final_score = round(min(score, 1.0), 2)
    print(f"Total score: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# ------------------ Script execution ------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/on_slide_241_the_heading_disappears_because_the_background_is_still_the_default_light_grey_i_want_th_golden.pptx"
    verify_task(FILE_PATH)
