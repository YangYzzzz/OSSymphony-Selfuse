"""
FINAL REWARD SCRIPT - SUCCESS
Task: I’ve got three photos all labeled “Picture 1” on Slide 3, Slide 4, and Slide 6 in my Impress deck. Could you walk me through the steps to resize each one so the height is exactly 20 cm on Slide 3, 30 cm on Slide 4, and 25 cm on Slide 6, while leaving the aspect ratio locked so they don’t get distorted?
Generated: 2025-09-10 13:26:07
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# -----------------------------------------------------------------------------
# Reward Script for Verifying Image Resizing in an Impress (PPTX) Deck
# -----------------------------------------------------------------------------
# Task Recap:
#   • Slide 3  → picture height must be 20 cm
#   • Slide 4  → picture height must be 30 cm
#   • Slide 6  → picture height must be 25 cm
#   • Aspect ratio must remain locked (width/height constant)
#
# Scoring (progressive):
#   • 0.3 pts per slide whose picture height matches the target (±0.25 cm)
#   • 0.1 pts if all verified pictures keep a consistent aspect ratio (±0.01)
#   → Maximum score = 1.0
# -----------------------------------------------------------------------------

FILE_PATH = (
    "/home/user/ive_got_three_photos_all_labeled_picture_1_on_slide_3_"
    "slide_4_and_slide_6_in_my_impress_deck_could_y_golden.pptx"
)

# EMU (English Metric Unit) conversion — 914400 EMU per inch; 1 inch = 2.54 cm
EMU_PER_CM = 914400 / 2.54  # 360 000 EMU ≈ 1 cm
TOLERANCE_CM = 0.25         # Allow ±0.25 cm for minor rounding differences

# Expected heights (cm) keyed by slide number (1-indexed)
EXPECTED_HEIGHTS = {3: 20, 4: 30, 6: 25}


def _picture_shapes(slide):
    """Return all picture shapes on a given slide."""
    return [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]


def _verify_picture_height(slide, expected_cm):
    """Check first picture height against expectation. Return tuple(result, h_cm, w_cm)."""
    pics = _picture_shapes(slide)
    if not pics:
        print("  ✗ No picture found on this slide")
        return False, None, None

    pic = pics[0]  # Assumes the relevant image is the first (only) picture
    h_cm = pic.height / EMU_PER_CM
    w_cm = pic.width  / EMU_PER_CM

    ok = abs(h_cm - expected_cm) <= TOLERANCE_CM
    status = "OK" if ok else "NOT OK"
    print(f"  Found picture: height={h_cm:.2f} cm, width={w_cm:.2f} cm (expected {expected_cm} cm) → {status}")
    return ok, h_cm, w_cm


def verify_task(pptx_path):
    """Main verification routine returning a score in [0.0, 1.0]."""
    print(f"Verifying presentation: {pptx_path}\n")

    # Guard: file must exist and loadable
    if not os.path.exists(pptx_path):
        print("✗ Presentation file not found")
        return 0.0
    try:
        prs = Presentation(pptx_path)
    except Exception as exc:
        print(f"✗ Failed to load presentation: {exc}")
        return 0.0

    total_score = 0.0
    aspect_ratios = []  # Collect width/height ratios of successfully-verified pictures

    # 1. Per-slide height verification (0.3 pts each)
    for slide_num, target_cm in EXPECTED_HEIGHTS.items():
        if slide_num > len(prs.slides):
            print(f"✗ Slide {slide_num} is missing from the deck")
            continue

        print(f"Slide {slide_num}:")
        ok, h_cm, w_cm = _verify_picture_height(prs.slides[slide_num - 1], target_cm)
        if ok:
            total_score += 0.3
            if h_cm:  # Avoid division by zero
                aspect_ratios.append(w_cm / h_cm)
        print()

    # 2. Aspect-ratio consistency check (0.1 pts)
    if len(aspect_ratios) == len(EXPECTED_HEIGHTS):
        avg_ratio = sum(aspect_ratios) / len(aspect_ratios)
        consistent = all(abs(r - avg_ratio) <= 0.01 for r in aspect_ratios)
        if consistent:
            total_score += 0.1
            print("✓ Aspect ratio locked across all resized images (0.1 pts)")
        else:
            print("✗ Aspect ratio mismatch detected (0 pts)")
    else:
        print("✗ Aspect ratio check skipped due to earlier failures (0 pts)")

    # Final score (cap at 1.0 & tidy float precision)
    final_score = round(min(total_score, 1.0), 6)
    if abs(final_score - 1.0) < 1e-6:
        final_score = 1.0

    print(f"\nTotal score: {final_score:.2f}")
    return final_score


if __name__ == "__main__":
    reward = verify_task(FILE_PATH)
    print(f"REWARD: {reward}")

