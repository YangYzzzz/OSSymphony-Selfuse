"""
Initial Setup: Draft_Slides.pptx with inconsistent formatting
Task ID: impress_wf_003
Domain: libreoffice_impress

Creates 6 slides with mixed backgrounds, varied title fonts/sizes, no slide numbers.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_003'
OUTPUT = f'{WORKDIR}/Draft_Slides.pptx'
DESKTOP = f'{WORKDIR}/Desktop/Draft_Slides.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_background(slide, r, g, b):
    """Set solid background color on a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def add_title_textbox(slide, prs, text, font_name, font_size_pt, bold, color_rgb):
    """Add a title-style text box at the top of the slide."""
    left = Inches(0.5)
    top = Inches(0.3)
    width = prs.slide_width - Inches(1.0)
    height = Inches(1.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color_rgb)
    return txBox


def add_body_textbox(slide, prs, lines, top_inches=1.8):
    """Add body text content to a slide."""
    left = Inches(0.5)
    top = Inches(top_inches)
    width = prs.slide_width - Inches(1.0)
    height = Inches(4.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        run = p.runs[0]
        run.font.name = "Calibri"
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    return txBox


def create_initial():
    prs = Presentation()
    # Standard widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # Blank

    # --- Slide 1: Title slide - light blue background, Times New Roman title ---
    s1 = prs.slides.add_slide(blank_layout)
    set_background(s1, 0xD6, 0xEA, 0xF8)  # light blue
    add_title_textbox(s1, prs, "Q4 2025 Strategic Review", "Times New Roman", 36, True, (0x1A, 0x5C, 0x8E))
    add_body_textbox(s1, prs, [
        "Prepared by: Strategy & Operations Division",
        "Date: December 15, 2025",
        "Classification: Internal Use Only",
    ])

    # --- Slide 2: White background, Calibri title ---
    s2 = prs.slides.add_slide(blank_layout)
    set_background(s2, 0xFF, 0xFF, 0xFF)  # white
    add_title_textbox(s2, prs, "Revenue Performance Summary", "Calibri", 28, False, (0x00, 0x00, 0x00))
    add_body_textbox(s2, prs, [
        "Total Revenue: $14.2M (+12% YoY)",
        "Recurring Revenue: $9.8M (69% of total)",
        "New Customer Revenue: $3.1M",
        "Enterprise Segment: $8.5M",
        "SMB Segment: $5.7M",
    ])

    # --- Slide 3: Light blue-gray background, Comic Sans title ---
    s3 = prs.slides.add_slide(blank_layout)
    set_background(s3, 0xE8, 0xF0, 0xFE)  # light blue-gray
    add_title_textbox(s3, prs, "Customer Acquisition Metrics", "Comic Sans MS", 30, True, (0x2E, 0x4A, 0x62))
    add_body_textbox(s3, prs, [
        "New Enterprise Clients: 23 (target: 20)",
        "Average Deal Size: $185,000",
        "Sales Cycle: 47 days (down from 62)",
        "Win Rate: 34% (up 8 pts)",
        "Pipeline Coverage: 3.2x",
    ])

    # --- Slide 4: Pale yellow background, Georgia title ---
    s4 = prs.slides.add_slide(blank_layout)
    set_background(s4, 0xFD, 0xF2, 0xCC)  # pale yellow
    add_title_textbox(s4, prs, "Product Development Roadmap", "Georgia", 34, False, (0x5D, 0x3F, 0x1A))
    add_body_textbox(s4, prs, [
        "Phase 1: API Gateway Redesign (Complete)",
        "Phase 2: ML Pipeline Integration (In Progress)",
        "Phase 3: Dashboard 2.0 Launch (Q1 2026)",
        "Phase 4: Mobile App Release (Q2 2026)",
        "Budget Utilization: 78% of allocated $2.4M",
    ])

    # --- Slide 5: Light green background, Impact title ---
    s5 = prs.slides.add_slide(blank_layout)
    set_background(s5, 0xD5, 0xF5, 0xE3)  # light green
    add_title_textbox(s5, prs, "Team Performance & Headcount", "Impact", 26, True, (0x1B, 0x4F, 0x28))
    add_body_textbox(s5, prs, [
        "Engineering: 42 FTEs (3 open roles)",
        "Sales: 18 FTEs (2 open roles)",
        "Customer Success: 12 FTEs (fully staffed)",
        "Average eNPS: 62 (industry avg: 40)",
        "Voluntary Attrition: 8.3% (target: <10%)",
    ])

    # --- Slide 6: Light coral background, Verdana title ---
    s6 = prs.slides.add_slide(blank_layout)
    set_background(s6, 0xFA, 0xDB, 0xD8)  # light coral
    add_title_textbox(s6, prs, "Key Risks and Mitigation Plans", "Verdana", 32, False, (0x78, 0x28, 0x1F))
    add_body_textbox(s6, prs, [
        "Risk 1: Supply chain delays - Mitigation: dual-source strategy",
        "Risk 2: Regulatory changes in EU - Mitigation: legal review Q1",
        "Risk 3: Key talent retention - Mitigation: equity refresh program",
        "Risk 4: Competitive pricing pressure - Mitigation: value-add bundles",
    ])

    # Save to home directory and Desktop
    prs.save(OUTPUT)
    os.makedirs(os.path.dirname(DESKTOP), exist_ok=True)
    import shutil
    shutil.copy2(OUTPUT, DESKTOP)
    print(f'Initial file created: {OUTPUT}')
    print(f'Desktop copy created: {DESKTOP}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{DESKTOP}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
