"""
Reward Script: Uniform formatting for Draft_Slides.pptx
Task ID: impress_wf_003
Domain: libreoffice_impress
Scoring:
  Component 1 (0.35): All 6 slides have white (#FFFFFF) background
  Component 2 (0.35): All titles are Arial 32pt bold #333333
  Component 3 (0.30): Slide numbers present in bottom-right of every slide
"""

import os

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_003'

# Known title texts from the presentation (first textbox on each slide)
EXPECTED_TITLES = [
    "Q4 2025 Strategic Review",
    "Revenue Performance Summary",
    "Customer Acquisition Metrics",
    "Product Development Roadmap",
    "Team Performance & Headcount",
    "Key Risks and Mitigation Plans",
]


def get_slide_bg_color(slide):
    """Get slide background color as hex string, or None."""
    try:
        fill = slide.background.fill
        if fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb).upper()
        elif fill.type == 5:  # inherited from master
            master_fill = slide.slide_layout.slide_master.background.fill
            if master_fill.type == 1:
                return str(master_fill.fore_color.rgb).upper()
    except Exception:
        pass
    return None


def find_title_shape(slide):
    """Find the title shape (first textbox with a known title text)."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text in EXPECTED_TITLES:
                return shape
    return None


def find_slide_number_shape(slide, expected_num, slide_width, slide_height):
    """Find a text shape containing the slide number in the bottom-right area."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text == str(expected_num):
            # Check if it's in the bottom-right quadrant
            # Bottom-right means: left > 50% of slide width, top > 50% of slide height
            center_x = shape.left + shape.width / 2
            center_y = shape.top + shape.height / 2
            in_right_half = center_x > slide_width * 0.5
            in_bottom_half = center_y > slide_height * 0.5
            if in_right_half and in_bottom_half:
                return shape
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)
    if len(slides) != 6:
        print(f"FAIL: Expected 6 slides, found {len(slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Component 1: All 6 slides have white (#FFFFFF) background (0.35 points)
    try:
        white_count = 0
        for i, slide in enumerate(slides):
            bg_color = get_slide_bg_color(slide)
            if bg_color == "FFFFFF":
                white_count += 1
            else:
                print(f"  Slide {i+1} background: {bg_color} (not white)")

        if white_count == 6:
            print(f"PASS: Component 1 -- All 6 slides have white background (0.35 pts)")
            total_score += 0.35
        elif white_count >= 4:
            # Partial credit only if majority changed (initial has only 1 white)
            partial = 0.35 * (white_count / 6)
            print(f"PARTIAL: Component 1 -- {white_count}/6 slides have white background ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- Only {white_count}/6 slides have white background")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: All titles are Arial 32pt bold #333333 (0.35 points)
    try:
        correct_title_count = 0
        for i, slide in enumerate(slides):
            title_shape = find_title_shape(slide)
            if title_shape is None:
                print(f"  Slide {i+1}: title shape not found")
                continue

            # Check runs of the title text
            all_runs_correct = True
            checked_any = False
            for para in title_shape.text_frame.paragraphs:
                for run in para.runs:
                    if not (run.text or "").strip():
                        continue
                    checked_any = True
                    issues = []

                    # Font name
                    if run.font.name != "Arial":
                        issues.append(f"font={run.font.name}")
                        all_runs_correct = False

                    # Font size: 32pt = 406400 EMU
                    if run.font.size != Pt(32):
                        issues.append(f"size={run.font.size}")
                        all_runs_correct = False

                    # Bold
                    if run.font.bold is not True:
                        issues.append(f"bold={run.font.bold}")
                        all_runs_correct = False

                    # Color #333333
                    try:
                        if run.font.color.type is not None:
                            color_hex = str(run.font.color.rgb).upper()
                            if color_hex != "333333":
                                issues.append(f"color={color_hex}")
                                all_runs_correct = False
                        else:
                            issues.append("color=inherited")
                            all_runs_correct = False
                    except Exception:
                        issues.append("color=error")
                        all_runs_correct = False

                    if issues:
                        print(f"  Slide {i+1} title issues: {', '.join(issues)}")

            if checked_any and all_runs_correct:
                correct_title_count += 1

        if correct_title_count == 6:
            print(f"PASS: Component 2 -- All 6 titles are Arial 32pt bold #333333 (0.35 pts)")
            total_score += 0.35
        elif correct_title_count > 0:
            partial = 0.35 * (correct_title_count / 6)
            print(f"PARTIAL: Component 2 -- {correct_title_count}/6 titles correct ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 -- No titles have correct formatting")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Slide numbers present in bottom-right of every slide (0.30 points)
    try:
        number_count = 0
        for i, slide in enumerate(slides):
            expected_num = i + 1
            num_shape = find_slide_number_shape(slide, expected_num, slide_width, slide_height)
            if num_shape is not None:
                number_count += 1
            else:
                print(f"  Slide {i+1}: slide number not found in bottom-right")

        if number_count == 6:
            print(f"PASS: Component 3 -- All 6 slides have slide numbers in bottom-right (0.30 pts)")
            total_score += 0.30
        elif number_count > 0:
            partial = 0.30 * (number_count / 6)
            print(f"PARTIAL: Component 3 -- {number_count}/6 slides have slide numbers ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 -- No slide numbers found")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
# The file is on the Desktop
file_path = f'{WORKDIR}/Desktop/Draft_Slides.pptx'
if not os.path.exists(file_path):
    # Fallback to home dir
    file_path = f'{WORKDIR}/Draft_Slides.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
