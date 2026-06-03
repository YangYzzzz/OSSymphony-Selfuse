"""
Reward Script: Verify three additional slide layouts in master slide
Task ID: impress_gf2_032
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): 'Title Only' layout exists with correct name
  Component 2 (0.15): 'Title Only' has a single title placeholder ~60% slide height
  Component 3 (0.15): 'Blank' layout exists with no placeholders
  Component 4 (0.20): 'Full Bleed Image' layout exists with full-slide picture placeholder
  Component 5 (0.15): 'Full Bleed Image' has overlay rectangle covering full slide
  Component 6 (0.15): Overlay rectangle is semi-transparent (40% opacity) and dark
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_032'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Build a map of layout names to layout objects
    layout_map = {}
    for layout in prs.slide_layouts:
        layout_map[layout.name] = layout

    print(f"INFO: Found {len(prs.slide_layouts)} layouts: {list(layout_map.keys())}")

    # We need to check that the 3 NEW layouts exist beyond the original 2 ('Default', 'Title, Content')
    # The initial file has only 'Default' and 'Title, Content'

    # Component 1: 'Title Only' layout exists (0.20 points)
    try:
        if 'Title Only' in layout_map:
            # Verify it's not one of the original layouts
            # Check it has at least one placeholder of TITLE type
            lo = layout_map['Title Only']
            has_title_ph = False
            for ph in lo.placeholders:
                ph_type = ph.placeholder_format.type
                # TITLE type is 1
                if ph_type is not None and int(ph_type) == 1:
                    has_title_ph = True
                    break
            if has_title_ph:
                print(f"PASS: Component 1 — 'Title Only' layout exists with title placeholder (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 1 — 'Title Only' layout exists but has no title placeholder")
        else:
            print(f"FAIL: Component 1 — 'Title Only' layout not found in layouts")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: 'Title Only' title placeholder height ~60% of slide height (0.15 points)
    try:
        if 'Title Only' in layout_map:
            lo = layout_map['Title Only']
            title_ph = None
            for ph in lo.placeholders:
                ph_type = ph.placeholder_format.type
                if ph_type is not None and int(ph_type) == 1:
                    title_ph = ph
                    break
            if title_ph is not None:
                height_pct = title_ph.height / slide_height * 100
                # Allow tolerance: 50-70% (task says ~60%)
                if 50.0 <= height_pct <= 70.0:
                    print(f"PASS: Component 2 — Title placeholder height is {height_pct:.1f}% of slide (0.15 pts)")
                    total_score += 0.15
                else:
                    print(f"FAIL: Component 2 — Title placeholder height is {height_pct:.1f}% (expected ~60%)")
            else:
                print(f"FAIL: Component 2 — No title placeholder in 'Title Only' layout")
        else:
            print(f"FAIL: Component 2 — 'Title Only' layout not found")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: 'Blank' layout exists with no content placeholders (0.15 points)
    # Blank layout should have no placeholders (or only metadata placeholders like date/footer/slide number)
    try:
        if 'Blank' in layout_map:
            lo = layout_map['Blank']
            # Count non-metadata placeholders (not date, footer, slide number)
            content_phs = []
            for ph in lo.placeholders:
                ph_type = ph.placeholder_format.type
                # Metadata types: DATE (16), FOOTER (15), SLIDE_NUMBER (13)
                if ph_type is not None and int(ph_type) not in (13, 15, 16):
                    content_phs.append(ph)
            if len(content_phs) == 0:
                print(f"PASS: Component 3 — 'Blank' layout has no content placeholders (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 3 — 'Blank' layout has {len(content_phs)} content placeholders")
        else:
            print(f"FAIL: Component 3 — 'Blank' layout not found in layouts")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: 'Full Bleed Image' layout exists with full-slide picture placeholder (0.20 points)
    try:
        if 'Full Bleed Image' in layout_map:
            lo = layout_map['Full Bleed Image']
            pic_ph = None
            for ph in lo.placeholders:
                ph_type = ph.placeholder_format.type
                # PICTURE type is 18
                if ph_type is not None and int(ph_type) == 18:
                    pic_ph = ph
                    break
            if pic_ph is not None:
                # Check dimensions match full slide (with 1% tolerance)
                w_match = abs(pic_ph.width - slide_width) / slide_width < 0.01
                h_match = abs(pic_ph.height - slide_height) / slide_height < 0.01
                pos_ok = pic_ph.left <= slide_width * 0.01 and pic_ph.top <= slide_height * 0.01
                if w_match and h_match and pos_ok:
                    print(f"PASS: Component 4 — 'Full Bleed Image' has full-slide picture placeholder (0.20 pts)")
                    total_score += 0.20
                else:
                    print(f"FAIL: Component 4 — Picture placeholder dimensions don't match slide: "
                          f"w={pic_ph.width} vs {slide_width}, h={pic_ph.height} vs {slide_height}")
            else:
                print(f"FAIL: Component 4 — 'Full Bleed Image' layout has no picture placeholder")
        else:
            print(f"FAIL: Component 4 — 'Full Bleed Image' layout not found")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: 'Full Bleed Image' has an overlay rectangle covering full slide (0.15 points)
    try:
        if 'Full Bleed Image' in layout_map:
            lo = layout_map['Full Bleed Image']
            overlay_found = False
            for shape in lo.shapes:
                if shape.is_placeholder:
                    continue
                # Check if this shape covers the full slide (with tolerance)
                w_match = abs(shape.width - slide_width) / slide_width < 0.02
                h_match = abs(shape.height - slide_height) / slide_height < 0.02
                pos_ok = shape.left <= slide_width * 0.02 and shape.top <= slide_height * 0.02
                if w_match and h_match and pos_ok:
                    overlay_found = True
                    print(f"PASS: Component 5 — Full-slide overlay shape found: '{shape.name}' (0.15 pts)")
                    total_score += 0.15
                    break
            if not overlay_found:
                print(f"FAIL: Component 5 — No full-slide overlay rectangle found in 'Full Bleed Image'")
        else:
            print(f"FAIL: Component 5 — 'Full Bleed Image' layout not found")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Overlay rectangle is semi-transparent (~40% opacity) and dark (0.15 points)
    # Need to check XML for alpha value
    try:
        import re as re_mod
        if 'Full Bleed Image' in layout_map:
            # Find the layout XML file for Full Bleed Image by reading raw XML from ZIP
            alpha_found = False
            dark_fill_found = False
            with zipfile.ZipFile(file_path, 'r') as zf:
                layout_files = sorted([n for n in zf.namelist()
                                       if 'slideLayout' in n and n.endswith('.xml') and '_rels' not in n])
                for lf in layout_files:
                    with zf.open(lf) as f:
                        raw_content = f.read().decode('utf-8')
                        # Check if this is the Full Bleed Image layout
                        if 'Full Bleed Image' not in raw_content:
                            continue
                        # Verify it's the layout name (in cSld name attribute)
                        name_match = re_mod.search(r'cSld[^>]*name="Full Bleed Image"', raw_content)
                        if not name_match:
                            continue
                        # Search raw XML for alpha and color values
                        alpha_matches = re_mod.findall(r'alpha val="(\d+)"', raw_content)
                        srgb_matches = re_mod.findall(r'srgbClr val="([0-9A-Fa-f]{6})"', raw_content)

                        for a in alpha_matches:
                            alpha_val = int(a)
                            # 40% opacity = alpha val 40000 in OOXML (100000 = fully opaque)
                            # Allow range 30000-50000
                            if 30000 <= alpha_val <= 50000:
                                alpha_found = True

                        for clr in srgb_matches:
                            # Check if color is dark (R, G, B all < 128)
                            r_val = int(clr[0:2], 16)
                            g_val = int(clr[2:4], 16)
                            b_val = int(clr[4:6], 16)
                            if r_val < 128 and g_val < 128 and b_val < 128:
                                dark_fill_found = True

                        break

            if alpha_found and dark_fill_found:
                print(f"PASS: Component 6 — Overlay has semi-transparent (~40% opacity) dark fill (0.15 pts)")
                total_score += 0.15
            elif alpha_found:
                print(f"FAIL: Component 6 — Alpha found but fill color is not dark")
            elif dark_fill_found:
                print(f"FAIL: Component 6 — Dark fill found but no semi-transparent alpha (~40%)")
            else:
                print(f"FAIL: Component 6 — No semi-transparent dark overlay found in XML")
        else:
            print(f"FAIL: Component 6 — 'Full Bleed Image' layout not found")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
