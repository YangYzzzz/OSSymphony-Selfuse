"""
Initial Setup: Configure master slide with layout variants
Task ID: impress_gf2_032
Domain: libreoffice_impress

Creates a presentation with 18 slides using only 2 layouts ('Default' and 'Title, Content').
The master slide has no additional layout variants - those are to be added by the agent.
"""

import os
import shlex
import subprocess
import time
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_032'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard 16:9 slide size
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)

    # We'll use the default template which has multiple layouts.
    # We need to build a presentation that effectively uses only 2 layout types:
    # Layout 0 = Title Slide (we'll call it 'Default')
    # Layout 1 = Title and Content (we'll call it 'Title, Content')
    # The task says master has "only two layouts: 'Default' and 'Title, Content'"

    # Access the slide master and rename layouts
    slide_master = prs.slide_masters[0]

    # Get the two layouts we want to keep references to
    layout_default = prs.slide_layouts[0]    # Title Slide -> rename to 'Default'
    layout_title_content = prs.slide_layouts[1]  # Title and Content -> rename to 'Title, Content'

    # Rename these layouts
    layout_default.name = 'Default'
    layout_title_content.name = 'Title, Content'

    # Remove all other layouts from the master by manipulating XML
    # Keep only indices 0 and 1
    from lxml import etree
    sldLayoutIdLst = slide_master.element.find(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}sldLayoutIdLst'
    )
    if sldLayoutIdLst is not None:
        layout_ids = list(sldLayoutIdLst)
        # Keep only first two
        for lid in layout_ids[2:]:
            sldLayoutIdLst.remove(lid)

    # Slide content data for a "Flexible Template" presentation
    slide_data = [
        # (layout_idx, title, content)
        (0, "Flexible Template", "Q1 2026 Strategy Overview"),
        (1, "Executive Summary", "Our organization achieved 23% revenue growth in FY2025, surpassing targets by 8 percentage points. Key drivers included expansion into APAC markets, launch of the Aurora product line, and strategic partnerships with TechVista and Meridian Corp."),
        (1, "Revenue Breakdown", "North America: $12.4M (+15%)\nEurope: $8.7M (+28%)\nAPAC: $5.2M (+45%)\nLatin America: $2.1M (+12%)\nTotal: $28.4M"),
        (1, "Product Performance", "Aurora Suite: $9.8M (35% of revenue)\nNexus Platform: $7.2M (25%)\nVelocity Tools: $6.1M (21%)\nCustom Solutions: $5.3M (19%)"),
        (0, "Strategic Priorities", "Focus Areas for 2026"),
        (1, "Market Expansion", "Target markets for Q2-Q4 2026:\n- Southeast Asia (Vietnam, Thailand, Indonesia)\n- Eastern Europe (Poland, Czech Republic)\n- Middle East (UAE, Saudi Arabia)\nProjected incremental revenue: $4.8M"),
        (1, "Technology Roadmap", "Phase 1 (Q2): AI-powered analytics dashboard\nPhase 2 (Q3): Real-time collaboration features\nPhase 3 (Q4): Mobile-first redesign\nTotal R&D investment: $3.2M"),
        (1, "Team Growth Plan", "Engineering: +15 headcount (ML specialists, backend)\nSales: +8 headcount (APAC regional managers)\nCustomer Success: +5 headcount\nDesign: +3 headcount\nTotal new hires: 31"),
        (0, "Financial Outlook", "Projections & Targets"),
        (1, "Q2 2026 Targets", "Revenue target: $8.2M\nGross margin: 72%\nCustomer acquisition: 45 enterprise accounts\nChurn rate target: <2.5%\nNPS target: >65"),
        (1, "Annual Forecast", "FY2026 Revenue: $34.5M (+21%)\nEBITDA margin: 18%\nFree cash flow: $4.1M\nHeadcount EOY: 215"),
        (1, "Investment Priorities", "R&D: $3.2M (acceleration of Aurora AI)\nMarketing: $2.8M (brand awareness in new markets)\nInfrastructure: $1.5M (cloud migration Phase 2)\nTotal capex: $7.5M"),
        (0, "Customer Insights", "Voice of the Customer"),
        (1, "Satisfaction Metrics", "Overall CSAT: 4.2/5.0\nProduct quality: 4.5/5.0\nSupport responsiveness: 3.8/5.0\nOnboarding experience: 4.1/5.0\nValue for money: 4.3/5.0"),
        (1, "Key Accounts Update", "TechVista Corp: Renewed 3-year contract ($2.4M)\nMeridian Industries: Expanding to 500 seats\nPacific Health Group: Pilot started March 2026\nAtlas Financial: RFP submitted, decision pending"),
        (1, "Competitive Landscape", "Primary competitors:\n- Streamline Pro (market share: 22%)\n- DataFlow Suite (market share: 18%)\n- Our position (market share: 15%)\nDifferentiators: AI capabilities, integration depth, customer success model"),
        (0, "Next Steps", "Action Items & Timeline"),
        (1, "Immediate Actions", "Week 1-2: Finalize APAC hiring plan\nWeek 3-4: Launch Aurora AI beta program\nWeek 5-6: Complete Q2 marketing campaign briefs\nWeek 7-8: Board presentation preparation\nOwners: Sarah Chen (APAC), Marcus Lee (Product), Diana Reyes (Marketing)"),
    ]

    for layout_idx, title_text, content_text in slide_data:
        if layout_idx == 0:
            slide = prs.slides.add_slide(layout_default)
            if slide.shapes.title:
                slide.shapes.title.text = title_text
            # Add subtitle to placeholder 1 if available
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    ph.text = content_text
                    break
        else:
            slide = prs.slides.add_slide(layout_title_content)
            if slide.shapes.title:
                slide.shapes.title.text = title_text
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    ph.text = content_text
                    break

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Number of slides: {len(prs.slides)}')
    print(f'Number of layouts: {len(prs.slide_layouts)}')
    for i, layout in enumerate(prs.slide_layouts):
        print(f'  Layout {i}: {layout.name}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
