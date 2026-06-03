"""
Reward Script: Interactive quiz slide with multiple-choice question and hyperlinks
Task ID: impress_ps_007
Domain: libreoffice_impress
Scoring:
  Component 1: Question text on slide 8 (0.2 pts)
  Component 2: Four answer buttons with correct labels (0.3 pts)
  Component 3: Correct answer B hyperlinks to slide 9 (0.2 pts)
  Component 4: Wrong answers A/C/D hyperlink to slide 10 (0.2 pts)
  Component 5: Button styling - colored fill + white text (0.1 pts)
"""

import os
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_007'


def get_hyperlink_target(slide, shape):
    """Extract hyperlink target slide filename from a shape's XML."""
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    elem = shape._element
    hlinks = elem.findall('.//{%s}hlinkClick' % ns_a)
    targets = []
    for h in hlinks:
        action = h.get('action', '')
        rid = h.get('{%s}id' % ns_r)
        if rid and 'hlinksldjump' in action:
            try:
                rel = slide.part.rels[rid]
                targets.append(rel.target_ref)
            except Exception:
                pass
    return targets


def get_shape_text(shape):
    """Get all text from a shape's text frame."""
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text.strip()


def get_shape_fill_rgb(shape):
    """Get solid fill color as hex string, or None."""
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def get_text_color_rgb(shape):
    """Get the font color of the first run in the shape, or None."""
    if not shape.has_text_frame:
        return None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            try:
                if run.font.color.type is not None:
                    return str(run.font.color.rgb)
            except Exception:
                pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Must have at least 10 slides
    if len(prs.slides) < 10:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 10")
        print("REWARD: 0.0")
        return 0.0

    slide8 = prs.slides[7]  # 0-indexed
    shapes = list(slide8.shapes)

    # Identify question shape and button shapes
    question_shape = None
    button_shapes = {}  # label letter -> shape

    expected_labels = {
        'A': 'Evacuate immediately',
        'B': 'Sound the alarm',
        'C': 'Try to fight the fire',
        'D': 'Call your manager',
    }

    for shape in shapes:
        text = get_shape_text(shape)
        if not text:
            continue

        # Check if this is one of the answer buttons
        matched_label = None
        for letter, label_text in expected_labels.items():
            # Check for patterns like "A. Evacuate immediately" or just the label
            if label_text.lower() in text.lower() or (letter + '.') in text:
                matched_label = letter
                break

        if matched_label:
            button_shapes[matched_label] = shape
        elif 'fire' in text.lower() and '?' in text:
            # This is the question text
            question_shape = shape

    # Component 1: Question text on slide 8 (0.2 points)
    try:
        if question_shape is not None:
            q_text = get_shape_text(question_shape)
            has_question = 'what should you do first' in q_text.lower() and 'fire' in q_text.lower()
            # Check bold - derive from actual font property
            is_bold = any(
                run.font.bold is True
                for para in (question_shape.text_frame.paragraphs if question_shape.has_text_frame else [])
                for run in para.runs
            )

            if has_question and is_bold:
                print(f"PASS: Component 1 - Question text found and bold: '{q_text}' (0.2 pts)")
                total_score += 0.2
            elif has_question:
                print(f"PARTIAL: Component 1 - Question text found but not bold: '{q_text}' (0.1 pts)")
                total_score += 0.1
            else:
                print(f"FAIL: Component 1 - Question text not matching expected: '{q_text}'")
        else:
            print("FAIL: Component 1 - No question text shape found on slide 8")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Four answer buttons with correct labels (0.3 points)
    # Award 0.075 per correct button found
    try:
        buttons_found = 0
        for letter in ['A', 'B', 'C', 'D']:
            if letter in button_shapes:
                shape = button_shapes[letter]
                text = get_shape_text(shape)
                exp = expected_labels[letter]
                if exp.lower() in text.lower():
                    buttons_found += 1
                    print(f"  Button {letter}: FOUND - '{text}'")
                else:
                    print(f"  Button {letter}: text mismatch - '{text}' vs expected containing '{exp}'")
            else:
                print(f"  Button {letter}: NOT FOUND on slide 8")

        pts = 0.3 * (buttons_found / 4.0)
        if buttons_found == 4:
            print(f"PASS: Component 2 - All 4 answer buttons found ({pts} pts)")
            total_score += pts
        elif buttons_found > 0:
            print(f"PARTIAL: Component 2 - {buttons_found}/4 buttons found ({pts:.3f} pts)")
            total_score += pts
        else:
            print(f"FAIL: Component 2 - No answer buttons found")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Correct answer B hyperlinks to slide 9 (0.2 points)
    try:
        if 'B' in button_shapes:
            targets = get_hyperlink_target(slide8, button_shapes['B'])
            # slide9.xml is the target for the correct answer
            links_to_slide9 = any('slide9' in t for t in targets)
            if links_to_slide9:
                print(f"PASS: Component 3 - Button B hyperlinks to slide 9 (correct feedback) (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 3 - Button B links to {targets}, expected slide9.xml")
        else:
            print("FAIL: Component 3 - Button B not found, cannot check hyperlink")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Wrong answers A/C/D hyperlink to slide 10 (0.2 points)
    # Award partial credit per correct wrong-answer link
    try:
        wrong_links_correct = 0
        for letter in ['A', 'C', 'D']:
            if letter in button_shapes:
                targets = get_hyperlink_target(slide8, button_shapes[letter])
                links_to_slide10 = any('slide10' in t for t in targets)
                if links_to_slide10:
                    wrong_links_correct += 1
                    print(f"  Button {letter}: correctly links to slide 10 (try again)")
                else:
                    print(f"  Button {letter}: links to {targets}, expected slide10.xml")
            else:
                print(f"  Button {letter}: not found, cannot check hyperlink")

        pts = 0.2 * (wrong_links_correct / 3.0)
        if wrong_links_correct == 3:
            print(f"PASS: Component 4 - All wrong answers link to slide 10 ({pts} pts)")
            total_score += pts
        elif wrong_links_correct > 0:
            print(f"PARTIAL: Component 4 - {wrong_links_correct}/3 wrong answers link correctly ({pts:.3f} pts)")
            total_score += pts
        else:
            print(f"FAIL: Component 4 - No wrong answer hyperlinks found")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Button styling - colored fill + white text (0.1 points)
    # Check that buttons have solid fill colors and white text
    try:
        styled_buttons = 0
        for letter in ['A', 'B', 'C', 'D']:
            if letter in button_shapes:
                shape = button_shapes[letter]
                fill_rgb = get_shape_fill_rgb(shape)
                text_rgb = get_text_color_rgb(shape)
                has_fill = fill_rgb is not None
                has_white_text = text_rgb is not None and text_rgb.upper() == 'FFFFFF'
                if has_fill and has_white_text:
                    styled_buttons += 1
                    print(f"  Button {letter}: fill={fill_rgb}, text_color={text_rgb} - styled correctly")
                else:
                    print(f"  Button {letter}: fill={fill_rgb}, text_color={text_rgb} - missing styling")

        if styled_buttons >= 3:
            print(f"PASS: Component 5 - {styled_buttons}/4 buttons have fill + white text (0.1 pts)")
            total_score += 0.1
        elif styled_buttons >= 1:
            pts = 0.1 * (styled_buttons / 4.0)
            print(f"PARTIAL: Component 5 - {styled_buttons}/4 buttons styled ({pts:.3f} pts)")
            total_score += pts
        else:
            print(f"FAIL: Component 5 - No buttons have proper styling")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
