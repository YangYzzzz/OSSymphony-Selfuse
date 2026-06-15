"""
Initial Setup: Safety Training presentation with empty quiz slide
Task ID: impress_ps_007
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_007'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_content_slide(prs, title_text, body_lines):
    """Helper to add a title+content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def add_title_only_slide(prs, title_text):
    """Helper to add a slide with only a title."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Safety Training Program"
    slide1.placeholders[1].text = "Workplace Health & Safety Compliance\nQ2 2025 Update"

    # --- Slide 2: Agenda ---
    add_content_slide(prs, "Training Agenda", [
        "Fire Safety Procedures",
        "Emergency Evacuation Routes",
        "First Aid Basics",
        "Hazard Identification",
        "Personal Protective Equipment (PPE)",
        "Incident Reporting Protocol",
    ])

    # --- Slide 3: Fire Safety Overview ---
    add_content_slide(prs, "Fire Safety Overview", [
        "Over 70,000 workplace fires occur annually in the US",
        "Average property damage exceeds $2.4 billion per year",
        "Most fires are preventable with proper training",
        "Every employee must know their role during a fire",
        "Regular drills reduce evacuation time by 40%",
    ])

    # --- Slide 4: Types of Fire Extinguishers ---
    add_content_slide(prs, "Types of Fire Extinguishers", [
        "Class A: Ordinary combustibles (wood, paper, cloth)",
        "Class B: Flammable liquids (gasoline, oil, grease)",
        "Class C: Electrical equipment fires",
        "Class D: Combustible metals (magnesium, titanium)",
        "Class K: Kitchen fires (cooking oils and fats)",
    ])

    # --- Slide 5: PASS Technique ---
    add_content_slide(prs, "The PASS Technique", [
        "P - Pull the pin from the extinguisher",
        "A - Aim the nozzle at the base of the fire",
        "S - Squeeze the handle to release the agent",
        "S - Sweep from side to side at the base",
        "Practice regularly during scheduled fire drills",
    ])

    # --- Slide 6: Emergency Evacuation Procedures ---
    add_content_slide(prs, "Emergency Evacuation Procedures", [
        "Know your nearest two exit routes at all times",
        "Do not use elevators during a fire emergency",
        "Assist colleagues with mobility limitations",
        "Meet at designated assembly points outside",
        "Account for all personnel using the buddy system",
    ])

    # --- Slide 7: Alarm Systems ---
    add_content_slide(prs, "Alarm Systems & Notifications", [
        "Building fire alarms are tested monthly on the first Monday",
        "Pull stations are located near all stairwell entrances",
        "PA system provides verbal instructions during emergencies",
        "Text/email alerts sent via the SafeAlert notification system",
        "Silent alarms connect directly to local fire department",
    ])

    # --- Slide 8: EMPTY (quiz slide placeholder) ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    # Intentionally empty - agent needs to add quiz content here

    # --- Slide 9: Correct Feedback ---
    slide9 = add_title_only_slide(prs, "Correct! Well done.")

    # --- Slide 10: Try Again Feedback ---
    slide10 = add_title_only_slide(prs, "Not quite - try again!")

    # --- Slide 11: First Aid Basics ---
    add_content_slide(prs, "First Aid Basics", [
        "Check the scene for safety before approaching",
        "Call emergency services (911) immediately",
        "Apply direct pressure to stop bleeding",
        "Keep the injured person calm and still",
        "Do not move someone with a suspected spinal injury",
    ])

    # --- Slide 12: Hazard Identification ---
    add_content_slide(prs, "Hazard Identification Checklist", [
        "Inspect work areas weekly for potential hazards",
        "Report frayed electrical cords immediately",
        "Ensure chemical storage follows SDS requirements",
        "Keep emergency exits clear of obstructions",
        "Verify that safety signage is visible and current",
    ])

    # --- Slide 13: PPE Requirements ---
    add_content_slide(prs, "Personal Protective Equipment", [
        "Hard hats required in all construction zones",
        "Safety glasses mandatory in laboratory areas",
        "Hearing protection above 85 dB exposure levels",
        "Steel-toed boots in warehouse and loading zones",
        "Chemical-resistant gloves for hazardous material handling",
    ])

    # --- Slide 14: Incident Reporting ---
    add_content_slide(prs, "Incident Reporting Protocol", [
        "All incidents must be reported within 24 hours",
        "Use Form IR-2025 available on the company intranet",
        "Near-misses should also be documented",
        "Supervisor must co-sign the incident report",
        "Safety committee reviews all reports monthly",
    ])

    # --- Slide 15: Summary & Next Steps ---
    add_content_slide(prs, "Summary & Next Steps", [
        "Complete the online safety assessment by March 31",
        "Attend the hands-on fire drill on April 5",
        "Review your department-specific safety procedures",
        "Submit any safety concerns to safety@company.com",
        "Certification cards will be issued within 2 weeks",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
