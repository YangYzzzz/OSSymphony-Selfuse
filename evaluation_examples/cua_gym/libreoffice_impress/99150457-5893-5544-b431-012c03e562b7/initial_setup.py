"""
Initial Setup: Create Training_Deck.pptx with 8 slides, no notes on slides 2/4/6, no transitions
Task ID: impress_gf4_003
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_003'
FILENAME = 'Training_Deck.pptx'
OUTPUT = f'{WORKDIR}/{FILENAME}'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.paragraphs[0].text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    add_title_slide(prs, "Employee Training Deck", "Quarterly Onboarding & Development Program\nHR Department - Q2 2025")

    # --- Slide 2: Company Overview (NO notes) ---
    add_content_slide(prs, "Company Overview", [
        "Founded in 2012 with headquarters in San Francisco",
        "Over 3,200 employees across 14 global offices",
        "Revenue grew 34% year-over-year to $890M in FY2024",
        "Named a Top 50 Workplace by Glassdoor three consecutive years",
        "Core verticals: Enterprise SaaS, Data Analytics, Cloud Infrastructure",
    ])

    # --- Slide 3: Our Mission & Values ---
    s3 = add_content_slide(prs, "Our Mission & Values", [
        "Mission: Empower organizations with intelligent data solutions",
        "Integrity: We operate transparently and ethically in every interaction",
        "Innovation: Continuous improvement drives our product roadmap",
        "Collaboration: Cross-functional teamwork is the foundation of success",
        "Customer Focus: Every decision starts with the end-user experience",
    ])
    # Slide 3 has notes (it is NOT in the "add notes" requirement)
    s3.notes_slide.notes_text_frame.text = "Review these values with the team during orientation."

    # --- Slide 4: Q1 2025 Performance (NO notes) ---
    add_content_slide(prs, "Q1 2025 Performance Highlights", [
        "Total revenue: $237M (+18% QoQ)",
        "New enterprise contracts: 42 signed in Q1",
        "Customer retention rate: 96.3%",
        "Net Promoter Score improved from 72 to 78",
        "R&D investment: $48M allocated to AI initiatives",
    ])

    # --- Slide 5: Team Structure ---
    s5 = add_content_slide(prs, "Team Structure & Leadership", [
        "CEO: Rebecca Martinez — 15 years in enterprise technology",
        "CTO: David Park — Former VP of Engineering at CloudScale",
        "VP of Sales: Aisha Patel — Built EMEA sales from $0 to $120M",
        "VP of Product: James Liu — Led product at two YC startups",
        "HR Director: Sofia Ramirez — Oversees 3,200+ employee programs",
    ])
    s5.notes_slide.notes_text_frame.text = "Introduce new hires to department heads during the first week."

    # --- Slide 6: Project Roadmap (NO notes) ---
    add_content_slide(prs, "2025 Project Roadmap", [
        "Q1: Launch AI-powered analytics dashboard (Completed)",
        "Q2: Roll out self-service data integration platform",
        "Q3: Expand European data center capacity by 40%",
        "Q4: Release enterprise security compliance suite v2.0",
        "Ongoing: Monthly sprint reviews and customer feedback cycles",
    ])

    # --- Slide 7: Key Metrics Dashboard ---
    s7 = add_content_slide(prs, "Key Metrics & KPIs", [
        "Monthly Active Users: 1.8M (target: 2.0M by Q4)",
        "Average Response Time: 142ms (SLA: <200ms)",
        "Uptime: 99.97% across all production services",
        "Employee Satisfaction Index: 4.3/5.0",
        "Training Completion Rate: 88% within first 30 days",
    ])
    s7.notes_slide.notes_text_frame.text = "These KPIs are reviewed in the monthly all-hands meeting."

    # --- Slide 8: Next Steps & Action Items ---
    s8 = add_content_slide(prs, "Next Steps & Action Items", [
        "Complete onboarding checklist by end of first week",
        "Schedule 1-on-1 meetings with your direct manager",
        "Review department-specific training modules on the LMS",
        "Join at least two cross-functional Slack channels",
        "Submit feedback on the onboarding experience by Day 30",
    ])
    s8.notes_slide.notes_text_frame.text = "Remind participants to bookmark the internal wiki for reference."

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
