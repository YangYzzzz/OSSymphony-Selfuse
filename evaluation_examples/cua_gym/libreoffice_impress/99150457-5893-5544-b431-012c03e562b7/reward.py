"""
Reward Script: Add speaker notes, apply Fade transitions, and export PDF with notes
Task ID: impress_gf4_003
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Speaker notes on slides 2, 4, 6 with >= 2 sentences each
  Component 2 (0.30): Fade transition on all 8 slides
  Component 3 (0.15): Auto-advance timing of 6000ms on all slides
  Component 4 (0.10): Transition duration ~0.8s on all slides
  Component 5 (0.15): PDF file exported to Desktop as Presentation_With_Notes.pdf
"""

import os
import re
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_003'


def persist_app_state(domain):
    """Try to save any unsaved LibreOffice state via Ctrl+S."""
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        import time
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def count_sentences(text):
    """Count sentences in text. A sentence ends with . ! or ?"""
    if not text or not text.strip():
        return 0
    # Split on sentence-ending punctuation followed by space or end
    sentences = re.split(r'[.!?]+(?:\s|$)', text.strip())
    # Filter out empty strings
    sentences = [s.strip() for s in sentences if s.strip()]
    return len(sentences)


def verify_task(pptx_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load presentation
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {pptx_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    if num_slides != 8:
        print(f"WARNING: Expected 8 slides, found {num_slides}")

    # Component 1: Speaker notes on slides 2, 4, 6 with >= 2 sentences each (0.30 points)
    # Initial state: slides 2, 4, 6 have empty notes. Golden state: they have >= 2 sentences.
    notes_score = 0.0
    target_slides = [2, 4, 6]  # 1-indexed
    for slide_num in target_slides:
        try:
            idx = slide_num - 1
            if idx >= num_slides:
                print(f"FAIL: Component 1 — Slide {slide_num} does not exist")
                continue
            slide = prs.slides[idx]
            try:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                notes_text = ""

            sentence_count = count_sentences(notes_text)
            if sentence_count >= 2:
                print(f"PASS: Component 1 — Slide {slide_num} has {sentence_count} sentences in notes (0.10 pts)")
                notes_score += 0.10
            else:
                print(f"FAIL: Component 1 — Slide {slide_num} notes have {sentence_count} sentences (need >= 2), text: {repr(notes_text[:80])}")
        except Exception as e:
            print(f"ERROR: Component 1 — Slide {slide_num}: {e}")

    total_score += notes_score
    print(f"Component 1 subtotal: {notes_score:.2f}/0.30")

    # Components 2-4: Transition checks via XML parsing
    ns_map = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    fade_count = 0
    adv_count = 0
    dur_count = 0

    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            for i in range(1, num_slides + 1):
                try:
                    with zf.open(f'ppt/slides/slide{i}.xml') as f:
                        root = ET.parse(f).getroot()
                        tr = root.find('.//p:transition', ns_map)
                        if tr is None:
                            print(f"FAIL: Slide {i} — no transition found")
                            continue

                        # Check for fade child element
                        fade_el = tr.find('.//p:fade', ns_map)
                        if fade_el is not None:
                            fade_count += 1
                        else:
                            # List what transition type IS present
                            children = [c.tag.split('}')[-1] if '}' in c.tag else c.tag for c in tr]
                            print(f"FAIL: Slide {i} — expected 'fade' transition, found children: {children}")

                        # Check advTm (auto-advance timing) = 6000ms
                        adv_tm = tr.get('advTm')
                        if adv_tm is not None and int(adv_tm) == 6000:
                            adv_count += 1
                        else:
                            print(f"FAIL: Slide {i} — advTm expected 6000, found {adv_tm}")

                        # Check transition duration (~0.8s)
                        # python-pptx uses 'spd' attribute: slow=1000, med=750, fast=500
                        # OR 'dur' attribute in milliseconds
                        dur_attr = tr.get('dur')
                        spd_attr = tr.get('spd')
                        if dur_attr is not None:
                            dur_ms = int(dur_attr)
                            # Accept 700-1000ms as ~0.8s
                            if 700 <= dur_ms <= 1000:
                                dur_count += 1
                            else:
                                print(f"FAIL: Slide {i} — transition dur={dur_ms}ms, expected ~800ms")
                        elif spd_attr in ('med', 'slow'):
                            # med ~750ms, slow ~1000ms; both are reasonable for "0.8s"
                            # 'med' is the standard mapping for ~0.8s in python-pptx
                            dur_count += 1
                        else:
                            print(f"FAIL: Slide {i} — no dur/spd for transition speed, spd={spd_attr}")

                except KeyError:
                    print(f"FAIL: Slide {i} XML not found in ZIP")
                except Exception as e:
                    print(f"ERROR: Slide {i} transition check: {e}")
    except Exception as e:
        print(f"ERROR: Cannot open ZIP: {e}")

    # Component 2: Fade transition on all slides (0.30 points)
    if fade_count == num_slides:
        print(f"PASS: Component 2 — All {num_slides} slides have Fade transition (0.30 pts)")
        total_score += 0.30
    elif fade_count > 0:
        partial = round(0.30 * fade_count / num_slides, 2)
        print(f"PARTIAL: Component 2 — {fade_count}/{num_slides} slides have Fade transition ({partial} pts)")
        total_score += partial
    else:
        print(f"FAIL: Component 2 — No slides have Fade transition")

    # Component 3: Auto-advance timing 6000ms (0.15 points)
    if adv_count == num_slides:
        print(f"PASS: Component 3 — All {num_slides} slides have advTm=6000ms (0.15 pts)")
        total_score += 0.15
    elif adv_count > 0:
        partial = round(0.15 * adv_count / num_slides, 2)
        print(f"PARTIAL: Component 3 — {adv_count}/{num_slides} slides have advTm=6000ms ({partial} pts)")
        total_score += partial
    else:
        print(f"FAIL: Component 3 — No slides have advTm=6000ms")

    # Component 4: Transition duration ~0.8s AND fade present (0.10 points)
    # Only score duration when the fade transition is also present, so this
    # doesn't award points on initial_env which has transition elements but no fade.
    if fade_count == num_slides and dur_count == num_slides:
        print(f"PASS: Component 4 — All {num_slides} slides have Fade + ~0.8s duration (0.10 pts)")
        total_score += 0.10
    elif fade_count > 0 and dur_count > 0:
        # Partial: count slides that have BOTH fade and acceptable duration
        both_count = min(fade_count, dur_count)
        partial = round(0.10 * both_count / num_slides, 2)
        print(f"PARTIAL: Component 4 — {both_count}/{num_slides} slides have Fade + ~0.8s duration ({partial} pts)")
        total_score += partial
    else:
        print(f"FAIL: Component 4 — Slides lack Fade transition with ~0.8s duration")

    # Component 5: PDF exported to Desktop (0.15 points)
    pdf_path = '/home/user/Desktop/Presentation_With_Notes.pdf'
    try:
        if os.path.exists(pdf_path):
            pdf_size = os.path.getsize(pdf_path)
            if pdf_size > 1000:
                print(f"PASS: Component 5 — PDF exists at {pdf_path}, size={pdf_size} bytes (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 5 — PDF exists but too small ({pdf_size} bytes), likely empty/corrupt")
        else:
            print(f"FAIL: Component 5 — PDF not found at {pdf_path}")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/Training_Deck.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
