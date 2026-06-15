"""
Reward Script: Insert author notes into LibreOffice Impress slides
Task ID: osworld_multi_apps_impress_notes_import_011
Domain: libreoffice_impress
Scoring:
  Component 1: Alice's notes in slides 1-5   — 0.4 points
  Component 2: Bob's notes in slides 6-10    — 0.3 points
  Component 3: Carol's notes in slides 11-15 — 0.3 points
  Total: 1.0

Ground truth (from context + VM exploration):
  - Slides 1-5:  Alice's notes (e.g. slide 1 starts with "Slide 1 — Introduction: ...")
  - Slides 6-10: Bob's notes   (e.g. slide 6 starts with "Slide 6 — Quantitative Results: ...")
  - Slides 11-15: Carol's notes (e.g. slide 11 starts with "Slide 11 — Theoretical Implications: ...")

All scores are earned only when notes are non-empty AND correspond to the correct author/content.
Initial state: all slides have empty notes — so all components return 0.0 on initial_env.
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_multi_apps_impress_notes_import_011'
PPTX_PATH = '/home/user/Desktop/Research_Collab.pptx'

# Expected note prefixes per slide — derived from the ground-truth docx content
# These unique prefixes ensure correct author content is verified (not just any text)
ALICE_EXPECTED_PREFIXES = [
    "Slide 1 — Introduction:",        # slide 1
    "Slide 2 — Literature Review:",   # slide 2
    "Slide 3 — Theoretical Framework:",  # slide 3
    "Slide 4 — Methodology:",          # slide 4
    "Slide 5 — Data Collection:",      # slide 5
]

BOB_EXPECTED_PREFIXES = [
    "Slide 6 — Quantitative Results:", # slide 6
    "Slide 7 — Hypothesis Testing:",   # slide 7
    "Slide 8 — Qualitative Findings:", # slide 8
    "Slide 9 — Triangulation:",        # slide 9
    "Slide 10 — Robustness:",          # slide 10
]

CAROL_EXPECTED_PREFIXES = [
    "Slide 11 — Theoretical Implications:", # slide 11
    "Slide 12 — Practical Implications:",   # slide 12
    "Slide 13 — Limitations:",              # slide 13
    "Slide 14 — Conclusion:",               # slide 14
    "Slide 15 — References:",               # slide 15
]


def get_slide_notes(slide):
    """Return notes text for a slide, stripped. Returns '' if no notes."""
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation — if this fails, nothing can be verified
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Verify correct slide count (15 slides expected)
    num_slides = len(prs.slides)
    if num_slides != 15:
        print(f"CRITICAL: Expected 15 slides, found {num_slides}. Cannot verify notes assignment.")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Alice's notes correctly inserted into slides 1-5 (0.4 points)
    # These notes should be present in golden_env but empty in initial_env
    try:
        alice_passed = 0
        for slide_idx, expected_prefix in enumerate(ALICE_EXPECTED_PREFIXES):
            slide = prs.slides[slide_idx]  # slides are 0-indexed
            notes_text = get_slide_notes(slide)
            if notes_text and notes_text.startswith(expected_prefix):
                alice_passed += 1
                print(f"PASS: Slide {slide_idx + 1} — Alice note starts with expected prefix ('{expected_prefix[:40]}...')")
            else:
                print(f"FAIL: Slide {slide_idx + 1} — expected Alice note starting with '{expected_prefix[:50]}', found: {repr(notes_text[:60])}")

        if alice_passed == 5:
            print(f"PASS: Component 1 — All 5 Alice slides have correct notes (0.4 pts)")
            total_score += 0.4
        elif alice_passed >= 3:
            partial = round(0.4 * alice_passed / 5, 2)
            print(f"PARTIAL: Component 1 — {alice_passed}/5 Alice slides correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — Only {alice_passed}/5 Alice slides have correct notes")
    except Exception as e:
        print(f"ERROR: Component 1 (Alice slides 1-5): {e}")

    # Component 2: Bob's notes correctly inserted into slides 6-10 (0.3 points)
    try:
        bob_passed = 0
        for i, expected_prefix in enumerate(BOB_EXPECTED_PREFIXES):
            slide_idx = 5 + i  # slides 6-10 are at indices 5-9
            slide = prs.slides[slide_idx]
            notes_text = get_slide_notes(slide)
            if notes_text and notes_text.startswith(expected_prefix):
                bob_passed += 1
                print(f"PASS: Slide {slide_idx + 1} — Bob note starts with expected prefix ('{expected_prefix[:40]}...')")
            else:
                print(f"FAIL: Slide {slide_idx + 1} — expected Bob note starting with '{expected_prefix[:50]}', found: {repr(notes_text[:60])}")

        if bob_passed == 5:
            print(f"PASS: Component 2 — All 5 Bob slides have correct notes (0.3 pts)")
            total_score += 0.3
        elif bob_passed >= 3:
            partial = round(0.3 * bob_passed / 5, 2)
            print(f"PARTIAL: Component 2 — {bob_passed}/5 Bob slides correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — Only {bob_passed}/5 Bob slides have correct notes")
    except Exception as e:
        print(f"ERROR: Component 2 (Bob slides 6-10): {e}")

    # Component 3: Carol's notes correctly inserted into slides 11-15 (0.3 points)
    try:
        carol_passed = 0
        for i, expected_prefix in enumerate(CAROL_EXPECTED_PREFIXES):
            slide_idx = 10 + i  # slides 11-15 are at indices 10-14
            slide = prs.slides[slide_idx]
            notes_text = get_slide_notes(slide)
            if notes_text and notes_text.startswith(expected_prefix):
                carol_passed += 1
                print(f"PASS: Slide {slide_idx + 1} — Carol note starts with expected prefix ('{expected_prefix[:40]}...')")
            else:
                print(f"FAIL: Slide {slide_idx + 1} — expected Carol note starting with '{expected_prefix[:50]}', found: {repr(notes_text[:60])}")

        if carol_passed == 5:
            print(f"PASS: Component 3 — All 5 Carol slides have correct notes (0.3 pts)")
            total_score += 0.3
        elif carol_passed >= 3:
            partial = round(0.3 * carol_passed / 5, 2)
            print(f"PARTIAL: Component 3 — {carol_passed}/5 Carol slides correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — Only {carol_passed}/5 Carol slides have correct notes")
    except Exception as e:
        print(f"ERROR: Component 3 (Carol slides 11-15): {e}")

    final_score = round(min(total_score, 1.0), 4)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Run verification against the canonical artifact path
if not os.path.exists(PPTX_PATH):
    print(f"File not found: {PPTX_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(PPTX_PATH)
