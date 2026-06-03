"""
Initial Setup: Multi-author research presentation with notes import task
Task ID: osworld_multi_apps_impress_notes_import_011
Domain: libreoffice_impress

Creates:
  - /home/user/Desktop/Research_Collab.pptx  (15 slides, NO notes)
  - /home/user/Desktop/collab_notes.docx     (author notes organized by heading)
  - /home/user/Desktop/slide_assignments.txt (slide-to-author mapping)
Opens the presentation in LibreOffice Impress.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Pt as DocxPt

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_impress_notes_import_011'
PPTX_PATH = f'{DESKTOP}/Research_Collab.pptx'
DOCX_PATH = f'{DESKTOP}/collab_notes.docx'
TXT_PATH  = f'{DESKTOP}/slide_assignments.txt'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


# ---------------------------------------------------------------------------
# Slide data: 15 research slides (NO notes in initial state)
# ---------------------------------------------------------------------------
SLIDE_DATA = [
    # (title, body_text)
    # Alice's slides: 1-5
    (
        "Introduction to Collaborative Research",
        "Overview of the multi-author study\nResearch objectives and scope\nTeam composition and roles",
    ),
    (
        "Literature Review: Foundational Works",
        "Key papers from 2010–2020\nGap analysis in existing studies\nFramework selection rationale",
    ),
    (
        "Theoretical Framework",
        "Core hypotheses H1–H3\nConceptual model diagram\nVariable definitions and operationalization",
    ),
    (
        "Research Methodology",
        "Mixed-methods design\nSample size: N = 450 participants\nData collection instruments",
    ),
    (
        "Data Collection Procedures",
        "Survey administration timeline\nInterview protocol summary\nEthics committee approval — Ref #ECB-2024-112",
    ),
    # Bob's slides: 6-10
    (
        "Quantitative Analysis Results",
        "Descriptive statistics table\nPearson correlation r = 0.78 (p < 0.001)\nRegression model: R² = 0.61",
    ),
    (
        "Hypothesis Testing",
        "H1 supported: β = 0.42, SE = 0.06\nH2 partially supported\nH3 rejected — see Appendix B",
    ),
    (
        "Qualitative Findings",
        "Thematic analysis of 30 interviews\nTheme 1: Collaboration barriers\nTheme 2: Technology enablers\nTheme 3: Organizational culture",
    ),
    (
        "Triangulation of Evidence",
        "Convergent validity confirmed\nDivergent cases examined\nIntegrated interpretation model",
    ),
    (
        "Statistical Robustness Checks",
        "Multicollinearity diagnostics: VIF < 2.5\nResidual normality (Shapiro-Wilk p = 0.34)\nBootstrap CI (95%): [0.31, 0.53]",
    ),
    # Carol's slides: 11-15
    (
        "Discussion: Theoretical Implications",
        "Extends Social Exchange Theory\nChallenges prior assumptions in Kim & Lee (2019)\nNovel contribution to collaborative cognition literature",
    ),
    (
        "Discussion: Practical Implications",
        "Guidelines for cross-functional teams\nPolicy recommendations for institutions\nTraining program design framework",
    ),
    (
        "Limitations and Future Research",
        "Sample limited to APAC region\nCross-sectional design — longitudinal study needed\nProposed future studies: longitudinal (3-year) cohort",
    ),
    (
        "Conclusion",
        "Research questions answered\nKey contributions summarized\nCall for replication studies",
    ),
    (
        "References & Acknowledgements",
        "Full bibliography available in supplementary materials\nFunded by NRF Grant #2024-CRS-007\nThanks to all 450 participants",
    ),
]


def create_pptx():
    """Create 15-slide Research_Collab.pptx with NO speaker notes."""
    prs = Presentation()
    # Use standard widescreen dimensions
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    layout_title   = prs.slide_layouts[0]  # Title Slide
    layout_content = prs.slide_layouts[1]  # Title, Content

    for idx, (title_text, body_text) in enumerate(SLIDE_DATA):
        if idx == 0:
            slide = prs.slides.add_slide(layout_title)
            slide.shapes.title.text = title_text
            try:
                slide.placeholders[1].text = "Multi-Author Research Collaboration"
            except Exception:
                pass
        else:
            slide = prs.slides.add_slide(layout_content)
            slide.shapes.title.text = title_text
            try:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                first = True
                for line in body_text.split('\n'):
                    if first:
                        tf.paragraphs[0].text = line
                        first = False
                    else:
                        p = tf.add_paragraph()
                        p.text = line
            except Exception:
                pass
        # Intentionally NO notes — task requires agent to insert them

    prs.save(PPTX_PATH)
    print(f'Created: {PPTX_PATH}')


# ---------------------------------------------------------------------------
# Author notes content (realistic research notes)
# ---------------------------------------------------------------------------
ALICE_NOTES = {
    1: (
        "Slide 1 — Introduction: Emphasize that this study was initiated in response "
        "to the 2022 NSF call for collaborative research proposals. "
        "Mention the three partner universities: Stanford, NUS, and UCL. "
        "Total project duration is 36 months."
    ),
    2: (
        "Slide 2 — Literature Review: Highlight Smith et al. (2018) as the cornerstone "
        "paper. Note the critical gap: no prior study combined social network analysis "
        "with knowledge management theory in a multi-institutional context. "
        "Cite the 47 papers reviewed."
    ),
    3: (
        "Slide 3 — Theoretical Framework: Walk the audience through the conceptual model "
        "slowly. H1 posits a positive relationship between trust and knowledge sharing. "
        "H2 examines the moderating role of organisational culture. "
        "H3 tests boundary conditions under remote-work settings."
    ),
    4: (
        "Slide 4 — Methodology: The mixed-methods design has two phases. "
        "Phase 1: online survey across all three universities (N=450). "
        "Phase 2: semi-structured interviews (n=30, purposive sampling). "
        "Instruments validated via pilot study with 25 respondents."
    ),
    5: (
        "Slide 5 — Data Collection: Survey deployed on Qualtrics, March–June 2024. "
        "Response rate: 73%. Interviews conducted via Zoom, average duration 48 minutes. "
        "All data anonymised; ethics approval obtained 14 Feb 2024."
    ),
}

BOB_NOTES = {
    6: (
        "Slide 6 — Quantitative Results: The Pearson r of 0.78 is remarkably strong for "
        "organisational research. Remind audience that N=450 gives substantial statistical "
        "power. The regression accounts for 8 control variables including firm size, "
        "industry sector, and geographic dispersion."
    ),
    7: (
        "Slide 7 — Hypothesis Testing: For H1, the standardised beta of 0.42 indicates "
        "that a 1-SD increase in inter-team trust predicts a 0.42-SD increase in knowledge "
        "sharing behaviour. H2 partial support: effect stronger in collectivist cultures. "
        "H3 rejection was surprising — see Appendix B for post-hoc analysis."
    ),
    8: (
        "Slide 8 — Qualitative Findings: Thirty interviews across all three sites. "
        "Theme 1 (Barriers): time-zone differences and language cited most frequently. "
        "Theme 2 (Enablers): shared digital workspaces rated highest utility. "
        "Theme 3 (Culture): organisations with flat hierarchies showed 2× higher sharing."
    ),
    9: (
        "Slide 9 — Triangulation: Convergent validity was assessed using a sequential "
        "explanatory design. The six divergent cases were re-interviewed; four resolved "
        "through richer contextual data. The integrated model uses path coefficients "
        "from quantitative phase, interpreted through qualitative narratives."
    ),
    10: (
        "Slide 10 — Robustness: All VIF values below 2.5 confirm no multicollinearity. "
        "Shapiro-Wilk p=0.34 means residuals are approximately normal. "
        "Bootstrap CIs (1,000 iterations) confirm the main effect is stable. "
        "Sensitivity analysis removing outliers (n=12) did not change conclusions."
    ),
}

CAROL_NOTES = {
    11: (
        "Slide 11 — Theoretical Implications: Our findings extend Social Exchange Theory "
        "by demonstrating that reciprocity norms operate differently across virtual and "
        "co-located teams. This directly challenges Kim & Lee (2019) who found no "
        "significant trust-sharing link in virtual contexts."
    ),
    12: (
        "Slide 12 — Practical Implications: Recommend three-tier intervention: "
        "(1) trust-building workshops in project kick-off weeks, "
        "(2) structured knowledge-sharing platforms with incentive mechanisms, "
        "(3) cultural competency training for international team leads. "
        "ROI estimated at 18% productivity gain based on case studies."
    ),
    13: (
        "Slide 13 — Limitations: Geographic scope is the main limitation — results may "
        "not generalise to EMEA or Americas. The cross-sectional design cannot establish "
        "causality; proposed 3-year longitudinal follow-up study is in ethics review. "
        "Also note that all surveys were self-reported."
    ),
    14: (
        "Slide 14 — Conclusion: We have answered all three research questions. "
        "Main contribution: first empirical evidence of trust-knowledge sharing link "
        "in a tri-institutional virtual collaboration. Practical tools are freely "
        "available at the project website (research-collab.org)."
    ),
    15: (
        "Slide 15 — References: Full bibliography is in the supplementary PDF shared "
        "via the conference portal. Grant acknowledgement: NRF #2024-CRS-007. "
        "Special thanks to the 450 participants and our research assistants "
        "Priya Sharma, David Wong, and Emma Fischer."
    ),
}


def create_docx():
    """Create collab_notes.docx with notes organised under author headings."""
    doc = Document()
    doc.add_heading('Research Collaboration — Author Notes', level=0)

    # Alice section
    doc.add_heading('Alice', level=1)
    doc.add_paragraph(
        'Notes for slides 1–5 (Introduction through Data Collection)'
    )
    for slide_num, note_text in ALICE_NOTES.items():
        doc.add_heading(f'Slide {slide_num}', level=2)
        doc.add_paragraph(note_text)

    # Bob section
    doc.add_heading('Bob', level=1)
    doc.add_paragraph(
        'Notes for slides 6–10 (Analysis and Results)'
    )
    for slide_num, note_text in BOB_NOTES.items():
        doc.add_heading(f'Slide {slide_num}', level=2)
        doc.add_paragraph(note_text)

    # Carol section
    doc.add_heading('Carol', level=1)
    doc.add_paragraph(
        'Notes for slides 11–15 (Discussion through References)'
    )
    for slide_num, note_text in CAROL_NOTES.items():
        doc.add_heading(f'Slide {slide_num}', level=2)
        doc.add_paragraph(note_text)

    doc.save(DOCX_PATH)
    print(f'Created: {DOCX_PATH}')


def create_txt():
    """Create slide_assignments.txt mapping slides to authors."""
    lines = [
        '# Slide Assignments — Research_Collab.pptx',
        '# Format: Slide <number>: <Author>',
        '',
        'Slide 1: Alice',
        'Slide 2: Alice',
        'Slide 3: Alice',
        'Slide 4: Alice',
        'Slide 5: Alice',
        'Slide 6: Bob',
        'Slide 7: Bob',
        'Slide 8: Bob',
        'Slide 9: Bob',
        'Slide 10: Bob',
        'Slide 11: Carol',
        'Slide 12: Carol',
        'Slide 13: Carol',
        'Slide 14: Carol',
        'Slide 15: Carol',
    ]
    os.makedirs(DESKTOP, exist_ok=True)
    with open(TXT_PATH, 'w') as f:
        f.write('\n'.join(lines) + '\n')
    print(f'Created: {TXT_PATH}')


def create_initial():
    os.makedirs(DESKTOP, exist_ok=True)
    create_pptx()
    create_docx()
    create_txt()

    # GUI-ready startup: open the presentation in LibreOffice Impress
    # and open the docx in LibreOffice Writer so the agent can read it
    launch_gui(f'libreoffice --impress "{PPTX_PATH}"', delay_sec=3.0)
    launch_gui(f'libreoffice --writer "{DOCX_PATH}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress and Writer with DISPLAY=:0')


create_initial()
