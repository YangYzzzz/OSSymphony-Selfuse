"""
Initial Setup: Create a 12-slide corporate presentation for PDF export security task.
Task ID: impress_el_025
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_el_025'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_items):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, item in enumerate(bullet_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
    return slide


def add_blank_slide_with_textbox(prs, title_text, body_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title textbox
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    # Body textbox
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8), Inches(4.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = body_text
    p2.font.size = Pt(16)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(prs, "Internal Only: Q4 2025 Strategic Review",
                    "Meridian Technologies Inc. | Confidential")

    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "Revenue grew 18.3% year-over-year to $247.6M",
        "EBITDA margin expanded to 24.7%, up from 21.2%",
        "Customer retention rate at 94.8% across all segments",
        "Three new enterprise contracts signed (>$10M each)",
        "R&D investment increased 22% to accelerate product roadmap",
    ])

    # Slide 3: Financial Overview
    add_content_slide(prs, "Financial Overview - Q4 2025", [
        "Total Revenue: $247.6M (vs. $209.3M Q4 2024)",
        "Gross Profit: $168.4M (68.0% margin)",
        "Operating Expenses: $112.8M (SG&A + R&D)",
        "Net Income: $41.2M (16.6% net margin)",
        "Free Cash Flow: $53.7M",
    ])

    # Slide 4: Revenue Breakdown
    add_content_slide(prs, "Revenue Breakdown by Segment", [
        "Enterprise Solutions: $142.3M (57.5%)",
        "Cloud Services: $61.9M (25.0%)",
        "Professional Services: $29.7M (12.0%)",
        "Licensing & Maintenance: $13.7M (5.5%)",
    ])

    # Slide 5: Customer Metrics
    add_content_slide(prs, "Customer Metrics & Retention", [
        "Total Active Customers: 2,847 (+312 net new)",
        "Enterprise Accounts (>$500K ARR): 186",
        "Net Revenue Retention: 118.4%",
        "Average Contract Value: $87.2K (up from $74.8K)",
        "Customer Satisfaction Score: 4.6/5.0",
    ])

    # Slide 6: Product Roadmap
    add_content_slide(prs, "Product Roadmap - 2026 Priorities", [
        "AI-powered analytics dashboard (GA: Q1 2026)",
        "Multi-cloud deployment automation (Beta: Q2 2026)",
        "Enhanced security compliance module (GA: Q2 2026)",
        "Mobile application redesign (GA: Q3 2026)",
        "Integration marketplace expansion (Ongoing)",
    ])

    # Slide 7: Competitive Landscape
    add_blank_slide_with_textbox(prs, "Competitive Landscape Analysis",
        "Our market position strengthened in Q4 with the acquisition of DataStream "
        "Analytics. Key competitors Nexus Corp and Vertex Systems both reported "
        "flat growth, while our enterprise segment expanded by 23%. The mid-market "
        "remains contested with aggressive pricing from new entrants. Our "
        "differentiation strategy focusing on integration capabilities and "
        "customer success continues to drive wins in competitive evaluations.")

    # Slide 8: Workforce & Talent
    add_content_slide(prs, "Workforce & Talent Initiatives", [
        "Total Headcount: 1,842 (up from 1,614)",
        "Engineering Team: 723 (39.2% of workforce)",
        "Voluntary Attrition: 8.2% (industry avg: 13.4%)",
        "Diversity Hiring: 47% of new hires from underrepresented groups",
        "Employee Engagement Score: 82/100",
    ])

    # Slide 9: Risk Assessment
    add_content_slide(prs, "Risk Assessment & Mitigation", [
        "Cybersecurity: Enhanced SOC 2 Type II compliance achieved",
        "Supply Chain: Diversified cloud provider dependencies",
        "Regulatory: GDPR and CCPA audit completed with zero findings",
        "Market: Hedging strategy for FX exposure on EU contracts",
        "Talent: Retention bonuses for critical engineering roles",
    ])

    # Slide 10: Capital Allocation
    add_content_slide(prs, "Capital Allocation Strategy", [
        "R&D Investment: $48.2M (19.5% of revenue)",
        "Acquisitions Reserve: $75M earmarked for 2026",
        "Share Repurchase Program: $30M authorized",
        "Dividend: $0.42/share quarterly (2.1% yield)",
        "Debt Reduction: $25M principal payment scheduled Q1 2026",
    ])

    # Slide 11: Key Milestones - Next Quarter
    add_content_slide(prs, "Key Milestones - Q1 2026", [
        "Complete DataStream Analytics integration by Feb 28",
        "Launch AI analytics dashboard to beta customers",
        "Close Series D funding for subsidiary MeridianX",
        "Achieve FedRAMP Moderate authorization",
        "Open APAC regional headquarters in Singapore",
    ])

    # Slide 12: Appendix / Disclaimer
    add_blank_slide_with_textbox(prs, "Confidentiality Notice",
        "This document contains proprietary and confidential information belonging "
        "to Meridian Technologies Inc. It is intended solely for the use of "
        "authorized personnel within the organization. Unauthorized distribution, "
        "reproduction, or disclosure of this material is strictly prohibited. "
        "All financial projections are forward-looking statements subject to "
        "market conditions and other risks detailed in our SEC filings. "
        "Contact: investor.relations@meridiantech.com")

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=3.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
