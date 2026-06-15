"""
Reward Script: Change title text on slide 1 to bold 40pt Montserrat font in dark navy (#0A1F3D)
Task ID: impress_tct_071
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25) - Font name is Montserrat
  Component 2 (0.25) - Font size is 40pt (508000 EMU)
  Component 3 (0.25) - Font bold is True
  Component 4 (0.25) - Font color is #0A1F3D
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_071'


def persist_app_state():
    """Best-effort save for any unsaved LibreOffice state."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 1 slide
    if len(prs.slides) < 1:
        print("FAIL: Presentation has no slides")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]

    # Find the title shape on slide 1
    title_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name and 'Title' in shape.name:
            title_shape = shape
            break

    # Fallback: look for shape with the expected title text
    if title_shape is None:
        for shape in slide.shapes:
            if shape.has_text_frame:
                full_text = "".join(p.text for p in shape.text_frame.paragraphs)
                if 'Strategic Roadmap 2025' in full_text:
                    title_shape = shape
                    break

    if title_shape is None:
        print("FAIL: Could not find title shape on slide 1")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: text content should still be 'Strategic Roadmap 2025'
    full_text = "".join(p.text for p in title_shape.text_frame.paragraphs)
    if 'Strategic Roadmap 2025' not in full_text:
        print(f"FAIL: Title text changed. Expected 'Strategic Roadmap 2025', found: {repr(full_text)}")
        print("REWARD: 0.0")
        return 0.0

    # Get all non-empty runs from the title
    runs = []
    for para in title_shape.text_frame.paragraphs:
        for run in para.runs:
            if (run.text or "").strip():
                runs.append(run)

    if not runs:
        print("FAIL: No text runs found in title shape")
        print("REWARD: 0.0")
        return 0.0

    # We check properties across all non-empty runs; all must match for full credit per component
    # Component 1: Font name is Montserrat (0.25 points)
    try:
        all_montserrat = all(r.font.name == 'Montserrat' for r in runs)
        if all_montserrat:
            print(f"PASS: Component 1 - Font name is Montserrat (0.25 pts)")
            total_score += 0.25
        else:
            actual_names = [r.font.name for r in runs]
            print(f"FAIL: Component 1 - Expected font 'Montserrat', found: {actual_names}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Font size is 40pt = 508000 EMU (0.25 points)
    try:
        expected_size = 508000  # 40pt * 12700 EMU/pt
        all_correct_size = all(r.font.size == expected_size for r in runs)
        if all_correct_size:
            print(f"PASS: Component 2 - Font size is 40pt / 508000 EMU (0.25 pts)")
            total_score += 0.25
        else:
            actual_sizes = [r.font.size for r in runs]
            print(f"FAIL: Component 2 - Expected size 508000 EMU (40pt), found: {actual_sizes}")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Font bold is True (0.25 points)
    try:
        # bold can be True, False, or None (inherit). Only True counts.
        all_bold = all(r.font.bold is True for r in runs)
        if all_bold:
            print(f"PASS: Component 3 - Font bold is True (0.25 pts)")
            total_score += 0.25
        else:
            actual_bold = [r.font.bold for r in runs]
            print(f"FAIL: Component 3 - Expected bold=True, found: {actual_bold}")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Font color is #0A1F3D (0.25 points)
    try:
        actual_colors = []
        for r in runs:
            try:
                if r.font.color.type is not None:
                    actual_colors.append(str(r.font.color.rgb).upper())
                else:
                    actual_colors.append('None/theme')
            except Exception:
                actual_colors.append('error')

        if len(actual_colors) > 0 and all(c == '0A1F3D' for c in actual_colors):
            print(f"PASS: Component 4 - Font color is #0A1F3D (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 4 - Expected color #0A1F3D, found: {actual_colors}")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
