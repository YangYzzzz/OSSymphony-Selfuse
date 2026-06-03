"""
Initial Setup: Change the title text on slide 1 to bold 40pt Montserrat font in dark navy (#0A1F3D).
Task ID: impress_tct_071
Domain: libreoffice_impress

Creates an 8-slide corporate presentation with slide 1 title in 28pt Liberation Sans, black, regular weight.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_071'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title1 = slide1.shapes.title
    title1.text = ""
    tf = title1.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = "Strategic Roadmap 2025"
    run.font.name = "Liberation Sans"
    run.font.size = Pt(28)
    run.font.bold = False
    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    # Subtitle
    subtitle = slide1.placeholders[1]
    subtitle.text = ""
    sub_run = subtitle.text_frame.paragraphs[0].add_run()
    sub_run.text = "Confidential \u2014 Executive Leadership Review"
    sub_run.font.name = "Liberation Sans"
    sub_run.font.size = Pt(16)
    sub_run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    items = [
        "Company Performance Overview",
        "Market Analysis & Competitive Landscape",
        "Product Development Pipeline",
        "Financial Projections Q3\u2013Q4 2025",
        "Talent Acquisition & Retention Strategy",
        "Key Risks & Mitigation Plans",
        "Q&A and Next Steps",
    ]
    body2.paragraphs[0].text = items[0]
    for item in items[1:]:
        p = body2.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 3: Company Performance ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Company Performance Overview"
    body3 = slide3.placeholders[1].text_frame
    perf_items = [
        "Revenue: $142.3M (up 18% YoY)",
        "Gross Margin: 67.2% (target: 65%)",
        "Active Users: 3.8M monthly (up 24%)",
        "Net Promoter Score: 72 (industry avg: 54)",
        "Employee Headcount: 1,240 across 12 offices",
    ]
    body3.paragraphs[0].text = perf_items[0]
    for item in perf_items[1:]:
        p = body3.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 4: Market Analysis ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Market Analysis"
    body4 = slide4.placeholders[1].text_frame
    market_items = [
        "Total Addressable Market: $28.5B by 2027",
        "Our Market Share: 8.3% (up from 5.1%)",
        "Primary Competitors: Nextera Corp, BlueStar Analytics, Prism Solutions",
        "Emerging Trend: AI-driven automation adoption +40% in enterprise",
        "Customer Acquisition Cost decreased 12% via organic channels",
    ]
    body4.paragraphs[0].text = market_items[0]
    for item in market_items[1:]:
        p = body4.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 5: Product Pipeline ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Product Development Pipeline"
    body5 = slide5.placeholders[1].text_frame
    prod_items = [
        "v4.2 Launch: Advanced Analytics Dashboard (June 2025)",
        "v4.5 Beta: AI Copilot for Workflow Automation (Aug 2025)",
        "Mobile App Redesign: Phase 2 complete, Phase 3 in progress",
        "API Gateway v3: Enterprise SSO & rate limiting (Sept 2025)",
        "Infrastructure: Migration to multi-region cloud by Q4",
    ]
    body5.paragraphs[0].text = prod_items[0]
    for item in prod_items[1:]:
        p = body5.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 6: Financial Projections ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Financial Projections Q3\u2013Q4 2025"
    body6 = slide6.placeholders[1].text_frame
    fin_items = [
        "Q3 Revenue Target: $38.7M | Q4 Revenue Target: $42.1M",
        "Projected Annual Revenue: $162.4M (14% above initial forecast)",
        "EBITDA Margin Goal: 22% by year-end",
        "R&D Investment: 28% of revenue allocated to innovation",
        "Capital Expenditure: $8.2M for data center expansion",
        "Cash Reserves: $54.6M (runway > 18 months at current burn)",
    ]
    body6.paragraphs[0].text = fin_items[0]
    for item in fin_items[1:]:
        p = body6.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 7: Talent Strategy ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Talent Acquisition & Retention"
    body7 = slide7.placeholders[1].text_frame
    talent_items = [
        "Open Positions: 87 across Engineering, Sales, and Product",
        "Attrition Rate: 9.4% (industry benchmark: 13.2%)",
        "New DEI Initiative: Partnership with 6 universities",
        "Leadership Development Program: 34 participants in Cohort 3",
        "Remote Work Policy: Hybrid model with 3 flexible days/week",
    ]
    body7.paragraphs[0].text = talent_items[0]
    for item in talent_items[1:]:
        p = body7.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 8: Next Steps ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Next Steps & Action Items"
    body8 = slide8.placeholders[1].text_frame
    next_items = [
        "Board Review: Final strategy deck due June 15, 2025",
        "Engineering Sprint: v4.2 feature freeze by May 30",
        "Sales Kickoff: Regional workshops in Singapore, London, Austin",
        "Customer Advisory Board: Annual meeting scheduled July 10\u201311",
        "Follow-up: Bi-weekly executive sync every other Thursday",
    ]
    body8.paragraphs[0].text = next_items[0]
    for item in next_items[1:]:
        p = body8.add_paragraph()
        p.text = item
        p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
