"""
Initial Setup: Create a presentation with 4 slides. Slide 3 has title and text but no shapes.
Task ID: impress_ndo_037
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_037'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard 25.4cm x 19.05cm (default for python-pptx)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Quarterly Business Review"
    slide1.placeholders[1].text = "FY2025 Q4 Performance Summary\nPrepared by Strategy Division"

    # --- Slide 2: Content Slide with bullet points ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Key Highlights"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Revenue grew 18% year-over-year to $4.2B"
    p2 = body2.add_paragraph()
    p2.text = "Customer acquisition cost decreased by 12%"
    p2.level = 0
    p3 = body2.add_paragraph()
    p3.text = "Net promoter score improved from 62 to 71"
    p3.level = 0
    p4 = body2.add_paragraph()
    p4.text = "Three new product lines launched in APAC markets"
    p4.level = 0

    # --- Slide 3: Title + text, NO shapes ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Market Position Analysis"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Our market share expanded to 23.5% in the enterprise segment"
    bp1 = body3.add_paragraph()
    bp1.text = "Competitive landscape remains fragmented with top 5 players holding 61% share"
    bp1.level = 0
    bp2 = body3.add_paragraph()
    bp2.text = "Strategic partnerships with Meridian Corp and Vanguard Technologies are driving channel growth"
    bp2.level = 0
    bp3 = body3.add_paragraph()
    bp3.text = "Brand awareness survey results show 15-point increase in aided recall"
    bp3.level = 0

    # --- Slide 4: Another content slide ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Next Steps & Roadmap"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Finalize Q1 2026 budget allocation by January 15"
    bp4a = body4.add_paragraph()
    bp4a.text = "Launch Phase 2 of the digital transformation initiative"
    bp4a.level = 0
    bp4b = body4.add_paragraph()
    bp4b.text = "Expand sales team in European markets by 20%"
    bp4b.level = 0
    bp4c = body4.add_paragraph()
    bp4c.text = "Complete integration of recently acquired DataFlow Analytics platform"
    bp4c.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
