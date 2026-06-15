"""
Reward Script: Draw a circle on slide 3 centered at the middle of the slide.
              Make it 6cm in diameter with a red (#E74C3C) fill and no border.
Task ID: impress_ndo_037
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Oval shape exists on slide 3
  Component 2 (0.25): Circle is 6cm diameter (W == H == 2160000 EMU)
  Component 3 (0.25): Circle is centered on the slide
  Component 4 (0.15): Fill color is #E74C3C
  Component 5 (0.10): No border/outline
"""

import os

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_037'

# Tolerance for position/size comparisons (0.5%)
def approx_equal(a, b, tol=0.005):
    if a == b:
        return True
    if a == 0 or b == 0:
        return abs(a - b) < 10000  # ~0.03 cm tolerance for zero comparisons
    return abs(a - b) / max(abs(a), abs(b)) <= tol


def find_oval_on_slide(slide):
    """Find the first oval/circle auto-shape on the given slide."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                from pptx.enum.shapes import MSO_SHAPE
                if shape.auto_shape_type == MSO_SHAPE.OVAL:
                    return shape
            except Exception:
                pass
        # Also check by name pattern as fallback
        if 'oval' in shape.name.lower() or 'circle' in shape.name.lower():
            return shape
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]  # 0-indexed
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Find the oval shape on slide 3
    oval = find_oval_on_slide(slide3)

    # Component 1: Oval shape exists on slide 3 (0.25 points)
    # This is the core task-introduced change: initial has no oval on slide 3
    try:
        if oval is not None:
            print(f"PASS: Component 1 -- Oval shape found on slide 3: '{oval.name}' (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 -- No oval shape found on slide 3")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    if oval is None:
        # No oval means no further checks can pass
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: Circle dimensions are 6cm (W == H == 2160000 EMU) (0.25 points)
    # 6cm = 2160000 EMU (360000 EMU per cm)
    try:
        expected_size = 2160000  # 6cm in EMU
        w = oval.width
        h = oval.height
        is_circle = approx_equal(w, h, tol=0.02)  # W == H means circle
        is_6cm = approx_equal(w, expected_size, tol=0.02) and approx_equal(h, expected_size, tol=0.02)

        if is_circle and is_6cm:
            print(f"PASS: Component 2 -- Circle is 6cm diameter (W={w}, H={h}, expected ~{expected_size}) (0.25 pts)")
            total_score += 0.25
        elif is_circle:
            print(f"PARTIAL: Component 2 -- Shape is circular but wrong size (W={w}, H={h}, expected ~{expected_size})")
            total_score += 0.10
        elif is_6cm:
            print(f"PARTIAL: Component 2 -- Size correct but not circular (W={w}, H={h})")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 -- Expected 6cm circle (W=H={expected_size}), found W={w}, H={h}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Circle is centered on slide (0.25 points)
    # Center = (slide_width/2, slide_height/2)
    # Shape center = (left + width/2, top + height/2)
    try:
        shape_center_x = oval.left + oval.width // 2
        shape_center_y = oval.top + oval.height // 2
        slide_center_x = slide_width // 2
        slide_center_y = slide_height // 2

        centered_h = approx_equal(shape_center_x, slide_center_x, tol=0.02)
        centered_v = approx_equal(shape_center_y, slide_center_y, tol=0.02)

        if centered_h and centered_v:
            print(f"PASS: Component 3 -- Circle centered (shape center: {shape_center_x},{shape_center_y}, slide center: {slide_center_x},{slide_center_y}) (0.25 pts)")
            total_score += 0.25
        elif centered_h:
            print(f"PARTIAL: Component 3 -- Centered horizontally only (shape center Y={shape_center_y}, slide center Y={slide_center_y})")
            total_score += 0.10
        elif centered_v:
            print(f"PARTIAL: Component 3 -- Centered vertically only (shape center X={shape_center_x}, slide center X={slide_center_x})")
            total_score += 0.10
        else:
            print(f"FAIL: Component 3 -- Not centered (shape center: {shape_center_x},{shape_center_y}, slide center: {slide_center_x},{slide_center_y})")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Fill color is #E74C3C (0.15 points)
    try:
        fill = oval.fill
        if fill.type is not None and fill.type == 1:  # SOLID fill
            color_rgb = str(fill.fore_color.rgb).upper()
            if color_rgb == 'E74C3C':
                print(f"PASS: Component 4 -- Fill color is #E74C3C (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 4 -- Expected fill #E74C3C, found #{color_rgb}")
        else:
            print(f"FAIL: Component 4 -- Fill is not solid (type={fill.type})")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: No border/outline (0.10 points)
    # No border means: line fill type is None or BACKGROUND, or line width is 0
    try:
        line = oval.line
        line_fill_type = line.fill.type
        line_width = line.width

        # No outline: fill type is None (no line defined), or BACKGROUND (5), or width == 0
        no_border = (line_fill_type is None) or (line_fill_type == 5) or (line_width == 0 or line_width is None)

        if no_border:
            print(f"PASS: Component 5 -- No border (line fill type={line_fill_type}, width={line_width}) (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 5 -- Border detected (line fill type={line_fill_type}, width={line_width})")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
