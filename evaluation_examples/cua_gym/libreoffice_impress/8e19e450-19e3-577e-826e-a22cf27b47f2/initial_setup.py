"""
Initial Setup: Workshop deck with 6 slides (alternating instructional and exercise slides)
Task ID: osworld_impress_slide_duplication_reorder_010
Domain: libreoffice_impress

Creates a 6-slide workshop presentation:
- Slides 1, 3, 5: Instructional content slides
- Slides 2, 4, 6: Exercise slides (to be duplicated by agent)
"""

import os
import shlex
import subprocess
import time
import copy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_slide_duplication_reorder_010'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_slide_background(slide, r, g, b):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def add_title_text(slide, title_text, title_color=(0x1F, 0x39, 0x7D)):
    """Add a title text box at the top of a blank slide."""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.name = "Calibri"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*title_color)
    p.alignment = PP_ALIGN.LEFT


def add_body_text(slide, body_text, top_inch=1.7, color=(0x33, 0x33, 0x33)):
    """Add body text box below title."""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(top_inch), Inches(9), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_text):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = "Calibri"
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(*color)
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)


def create_initial():
    prs = Presentation()
    # Use blank layout for full control
    blank_layout = prs.slide_layouts[6]  # Blank

    # --- Slide 1: Instructional — Introduction to Data Analysis ---
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, 0xE8, 0xF0, 0xFE)  # light blue background
    add_title_text(slide1, "Module 1: Introduction to Data Analysis")
    add_body_text(slide1, [
        "Welcome to the Data Analysis Workshop",
        "",
        "In this module, you will learn:",
        "  • Understanding data types and structures",
        "  • Importing and cleaning data with Python",
        "  • Performing descriptive statistics",
        "  • Visualizing trends with matplotlib",
        "",
        "Prerequisites: Basic Python knowledge",
        "Duration: 45 minutes",
    ])

    # --- Slide 2: Exercise — Exercise 1: Data Cleaning ---
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, 0xFF, 0xF9, 0xE6)  # light yellow background for exercise
    add_title_text(slide2, "Exercise 1: Data Cleaning", title_color=(0x7B, 0x3F, 0x00))
    add_body_text(slide2, [
        "Task: Clean the provided sales dataset",
        "",
        "Instructions:",
        "  1. Load 'sales_q1_2024.csv' using pandas",
        "  2. Identify and remove rows with missing values",
        "  3. Convert the 'Date' column to datetime format",
        "  4. Remove duplicate entries based on 'OrderID'",
        "  5. Export the cleaned dataset as 'sales_clean.csv'",
        "",
        "Expected output: Dataset with 847 rows (from original 912)",
    ], color=(0x4A, 0x30, 0x00))

    # --- Slide 3: Instructional — Statistical Analysis ---
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, 0xE8, 0xF5, 0xE9)  # light green background
    add_title_text(slide3, "Module 2: Statistical Analysis Techniques", title_color=(0x1B, 0x5E, 0x20))
    add_body_text(slide3, [
        "Core Statistical Methods",
        "",
        "Descriptive Statistics:",
        "  • Mean, Median, Mode — central tendency measures",
        "  • Standard Deviation, Variance — spread measures",
        "  • Quartiles and Interquartile Range (IQR)",
        "",
        "Inferential Statistics:",
        "  • Hypothesis testing (t-test, chi-squared)",
        "  • Confidence intervals and p-values",
        "  • Correlation and regression analysis",
        "",
        "Python libraries: scipy.stats, statsmodels",
    ], color=(0x1B, 0x5E, 0x20))

    # --- Slide 4: Exercise — Exercise 2: Statistical Analysis ---
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, 0xFF, 0xF9, 0xE6)  # light yellow background for exercise
    add_title_text(slide4, "Exercise 2: Statistical Analysis", title_color=(0x7B, 0x3F, 0x00))
    add_body_text(slide4, [
        "Task: Analyze customer purchase patterns",
        "",
        "Instructions:",
        "  1. Load the cleaned dataset from Exercise 1",
        "  2. Calculate mean and median purchase amounts by region",
        "  3. Test whether the Northern and Southern regions differ",
        "     significantly (use two-sample t-test, α=0.05)",
        "  4. Compute the Pearson correlation between",
        "     'CustomerAge' and 'PurchaseAmount'",
        "  5. Report findings in a summary DataFrame",
        "",
        "Hint: Use scipy.stats.ttest_ind() for the t-test",
    ], color=(0x4A, 0x30, 0x00))

    # --- Slide 5: Instructional — Data Visualization ---
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, 0xF3, 0xE5, 0xF5)  # light purple background
    add_title_text(slide5, "Module 3: Data Visualization Best Practices", title_color=(0x4A, 0x14, 0x8C))
    add_body_text(slide5, [
        "Effective Data Visualization",
        "",
        "Chart Selection Guide:",
        "  • Bar charts — compare categories",
        "  • Line charts — show trends over time",
        "  • Scatter plots — reveal relationships",
        "  • Heatmaps — display correlation matrices",
        "  • Box plots — compare distributions",
        "",
        "Design Principles:",
        "  • Choose colors accessible to colorblind viewers",
        "  • Label axes clearly with units",
        "  • Include chart titles and legends",
        "  • Avoid 3D charts and unnecessary decoration",
    ], color=(0x4A, 0x14, 0x8C))

    # --- Slide 6: Exercise — Exercise 3: Data Visualization ---
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, 0xFF, 0xF9, 0xE6)  # light yellow background for exercise
    add_title_text(slide6, "Exercise 3: Data Visualization", title_color=(0x7B, 0x3F, 0x00))
    add_body_text(slide6, [
        "Task: Create a comprehensive visualization dashboard",
        "",
        "Instructions:",
        "  1. Generate a bar chart of total sales by product category",
        "  2. Plot a line chart of weekly revenue over Q1 2024",
        "  3. Create a scatter plot of CustomerAge vs PurchaseAmount",
        "     colored by region",
        "  4. Build a correlation heatmap for numeric columns",
        "  5. Arrange all 4 charts in a 2×2 subplot grid",
        "  6. Save as 'analysis_dashboard.png' (300 DPI)",
        "",
        "Bonus: Add trend line to the scatter plot using numpy.polyfit",
    ], color=(0x4A, 0x30, 0x00))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Slide count: {len(prs.slides)} (expected 6)')

    # GUI-ready startup: open the file in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
