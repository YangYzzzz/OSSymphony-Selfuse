"""
Reward Script: Duplicate slides 2, 4, 6 and insert each copy immediately after
               its original, resulting in 9-slide pattern: 1, 2, 2', 3, 4, 4', 5, 6, 6'
Task ID: osworld_impress_slide_duplication_reorder_010
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): Total slide count is exactly 9
  Component 2 (0.4): Correct ordering — slides at even positions 3,6,9 are exact
                     duplicates of slides 2,5,8 (exercise slides paired with copies)
  Component 3 (0.3): Original instructional slides remain at correct odd positions
                     (1->Module1, 4->Module2, 7->Module3) and maintain proper sequencing
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_slide_duplication_reorder_010'

# Known title texts from initial presentation (used to identify slide roles)
EXERCISE_TITLES = [
    'Exercise 1: Data Cleaning',
    'Exercise 2: Statistical Analysis',
    'Exercise 3: Data Visualization',
]
INSTRUCTION_TITLES = [
    'Module 1: Introduction to Data Analysis',
    'Module 2: Statistical Analysis Techniques',
    'Module 3: Data Visualization Best Practices',
]


def get_slide_all_text(slide):
    """Return a sorted tuple of all non-empty text strings from a slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return tuple(texts)


def get_slide_first_title(slide):
    """Return the first non-empty text from a slide (usually the title)."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    return t
    return ""


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print("CRITICAL: Cannot load file {}: {}".format(file_path, e))
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    print("INFO: Presentation has {} slides".format(num_slides))

    # Print all slide titles for diagnostics
    for i, slide in enumerate(prs.slides):
        title = get_slide_first_title(slide)
        print("  Slide {}: {}".format(i + 1, title[:60] if title else "(no text)"))

    # Component 1: Total slide count is exactly 9 (0.3 points)
    # In the initial env there are 6 slides; after the task there should be 9.
    try:
        if num_slides == 9:
            print("PASS: Component 1 — slide count is 9 (0.3 pts)")
            total_score += 0.3
        else:
            print("FAIL: Component 1 — expected 9 slides, found {}".format(num_slides))
    except Exception as e:
        print("ERROR: Component 1 — {}".format(e))

    # Component 2: Correct ordering — exercise pairs at positions (2,3), (5,6), (8,9) (0.4 points)
    # Pattern: 1, 2, 2', 3, 4, 4', 5, 6, 6'
    # i.e., slide 3 must be a duplicate of slide 2, slide 6 of slide 5, slide 9 of slide 8.
    # Verified by comparing all text content between paired slides.
    try:
        if num_slides >= 9:
            slides = list(prs.slides)
            # Expected pairs: (index 1, index 2), (index 4, index 5), (index 7, index 8) [0-based]
            pairs = [(1, 2), (4, 5), (7, 8)]
            pair_results = []
            for orig_idx, dup_idx in pairs:
                orig_texts = get_slide_all_text(slides[orig_idx])
                dup_texts = get_slide_all_text(slides[dup_idx])
                match = orig_texts == dup_texts and len(orig_texts) > 0
                pair_results.append(match)
                orig_title = get_slide_first_title(slides[orig_idx])
                dup_title = get_slide_first_title(slides[dup_idx])
                status = "MATCH" if match else "MISMATCH"
                print("  Pair (slides {}-{}): {} | '{}' vs '{}'".format(
                    orig_idx + 1, dup_idx + 1, status,
                    orig_title[:40], dup_title[:40]))

            if all(pair_results):
                print("PASS: Component 2 — all 3 exercise slide pairs match correctly (0.4 pts)")
                total_score += 0.4
            elif sum(pair_results) == 2:
                print("PARTIAL: Component 2 — 2/3 exercise slide pairs match (0.25 pts)")
                total_score += 0.25
            elif sum(pair_results) == 1:
                print("PARTIAL: Component 2 — 1/3 exercise slide pairs match (0.1 pts)")
                total_score += 0.1
            else:
                print("FAIL: Component 2 — no exercise slide pairs match")
        else:
            print("FAIL: Component 2 — insufficient slides ({}) to check pairs".format(num_slides))
    except Exception as e:
        print("ERROR: Component 2 — {}".format(e))

    # Component 3: Original slide sequencing preserved — instructional slides at correct
    # positions and exercise slides (originals) appear before their copies (0.3 points)
    # Expected positions (1-based):
    #   Slide 1: Module 1 (instructional)
    #   Slide 2: Exercise 1 (exercise)
    #   Slide 4: Module 2 (instructional)
    #   Slide 5: Exercise 2 (exercise)
    #   Slide 7: Module 3 (instructional)
    #   Slide 8: Exercise 3 (exercise)
    try:
        if num_slides >= 9:
            slides = list(prs.slides)
            # Check instructional slides at positions 1, 4, 7 (0-based: 0, 3, 6)
            instruction_checks = [
                (0, INSTRUCTION_TITLES[0]),   # Module 1 at slide 1
                (3, INSTRUCTION_TITLES[1]),   # Module 2 at slide 4
                (6, INSTRUCTION_TITLES[2]),   # Module 3 at slide 7
            ]
            # Check exercise slides (originals) at positions 2, 5, 8 (0-based: 1, 4, 7)
            exercise_checks = [
                (1, EXERCISE_TITLES[0]),   # Exercise 1 at slide 2
                (4, EXERCISE_TITLES[1]),   # Exercise 2 at slide 5
                (7, EXERCISE_TITLES[2]),   # Exercise 3 at slide 8
            ]
            all_checks = instruction_checks + exercise_checks
            passed_checks = 0
            for idx, expected_title in all_checks:
                actual_title = get_slide_first_title(slides[idx])
                if actual_title == expected_title:
                    passed_checks += 1
                    print("  PASS: Slide {} has correct title '{}'".format(idx + 1, expected_title[:50]))
                else:
                    print("  FAIL: Slide {} — expected '{}', found '{}'".format(
                        idx + 1, expected_title[:40], actual_title[:40]))

            if passed_checks == 6:
                print("PASS: Component 3 — all 6 slide positions correct (0.3 pts)")
                total_score += 0.3
            elif passed_checks >= 4:
                print("PARTIAL: Component 3 — {}/6 slide positions correct (0.15 pts)".format(passed_checks))
                total_score += 0.15
            elif passed_checks >= 2:
                print("PARTIAL: Component 3 — {}/6 slide positions correct (0.05 pts)".format(passed_checks))
                total_score += 0.05
            else:
                print("FAIL: Component 3 — only {}/6 slide positions correct".format(passed_checks))
        else:
            print("FAIL: Component 3 — insufficient slides to check ordering")
    except Exception as e:
        print("ERROR: Component 3 — {}".format(e))

    final_score = min(total_score, 1.0)
    print("\nScore: {}/1.0".format(total_score))
    print("REWARD: {}".format(final_score))
    return final_score


# Default: test against canonical artifact path
file_path = '{}/{}.pptx'.format(WORKDIR, TASK_ID)
if not os.path.exists(file_path):
    print("File not found: {}".format(file_path))
    print("REWARD: 0.0")
else:
    verify_task(file_path)
