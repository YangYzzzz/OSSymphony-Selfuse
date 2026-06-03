"""
FINAL REWARD SCRIPT - SUCCESS
Task: Hey, I'm working on a presentation and I need to start a new section called 'Appendix' beginning with slide 9. How can I set this up in LibreOffice Impress?
Generated: 2025-08-07 09:09:36
Status: success
Model: o4-mini
Total Steps: 4
"""

import os
from pptx import Presentation

def verify_appendix_section(file_path):
    """
    Verifies that a presentation file exists, has at least 9 slides,
    and that slide 9 contains the text 'Appendix'.
    Progressive scoring:
      - 0.2 for file existence
      - 0.2 for successful load
      - 0.2 for having >=9 slides
      - 0.4 for finding 'Appendix' on slide 9
    Returns a float between 0.0 and 1.0 and prints detailed debugging info.
    """
    print("Checking task completion: Section 'Appendix' starts at slide 9")
    score = 0.0
    max_score = 1.0

    # Requirement 1: file exists
    print(f"Verifying file existence: {file_path}")
    if os.path.exists(file_path):
        print("✓ File exists (0.2 pts)")
        score += 0.2
    else:
        print("✗ File not found")
        print(f"REWARD: {score}")
        return score

    # Requirement 2: load presentation
    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded with {len(prs.slides)} slides (0.2 pts)")
        score += 0.2
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        print(f"REWARD: {score}")
        return score

    # Requirement 3: slide count >=9
    slide_count = len(prs.slides)
    print(f"Verifying slide count: found {slide_count} slides, need >=9")
    if slide_count >= 9:
        print("✓ Slide count requirement met (0.2 pts)")
        score += 0.2
    else:
        print("✗ Not enough slides (<9)")

    # Requirement 4: slide 9 title or content contains 'Appendix'
    if slide_count >= 9:
        slide9 = prs.slides[8]
        title_found = False
        found_texts = []
        for shape in slide9.shapes:
            if hasattr(shape, 'text'):
                txt = shape.text.strip()
                if txt:
                    found_texts.append(txt)
                    if 'appendix' in txt.lower():
                        title_found = True
                        print(f"✓ Found text on slide 9: '{txt}'")
        if title_found:
            print("✓ Slide 9 contains 'Appendix' (0.4 pts)")
            score += 0.4
        else:
            print("✗ 'Appendix' not found on slide 9. Texts found:")
            for t in found_texts:
                print(f"  - {t}")

    # Finalize score
    final_score = min(score, max_score)
    print(f"Total score: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score

if __name__ == '__main__':
    # Golden file path for verification
    file_path = '/home/user/hey_im_working_on_a_presentation_and_i_need_to_start_a_new_section_called_appendix_beginning_with_sl.pptx'
    verify_appendix_section(file_path)
