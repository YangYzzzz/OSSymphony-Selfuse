"""
Initial Setup: Create a presentation with a grouped column chart (default colors)
Task ID: impress_tct_046
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_046'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Product Lines Performance"
    slide1.placeholders[1].text = "Q1-Q4 2025 Revenue Analysis"

    # --- Slide 2: Chart Slide (grouped column chart with 3 series, default colors) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Add a title text box at the top
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Quarterly Revenue by Product Line"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Chart data
    chart_data = CategoryChartData()
    chart_data.categories = ['Q1 2025', 'Q2 2025', 'Q3 2025', 'Q4 2025']
    chart_data.add_series('Enterprise Software', (245, 312, 287, 356))
    chart_data.add_series('Cloud Services', (189, 234, 267, 298))
    chart_data.add_series('Hardware Solutions', (156, 178, 145, 192))

    # Add grouped column chart
    chart_shape = slide2.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.5),
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = True

    # DO NOT set custom colors - leave defaults for the agent task

    # --- Slide 3: Analysis Slide ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Key Insights"
    p3.alignment = PP_ALIGN.LEFT
    run3 = p3.runs[0]
    run3.font.size = Pt(28)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

    body3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.0))
    tf_body3 = body3.text_frame
    tf_body3.word_wrap = True

    insights = [
        "Enterprise Software maintained the strongest revenue growth, reaching $356K in Q4.",
        "Cloud Services showed the most consistent quarter-over-quarter increase at 15.8% average.",
        "Hardware Solutions experienced a dip in Q3 but recovered strongly in Q4 with 32% growth.",
        "Total revenue across all product lines increased by 22% from Q1 to Q4.",
        "Cloud Services is projected to surpass Enterprise Software by Q2 2026."
    ]
    for i, insight in enumerate(insights):
        if i == 0:
            p = tf_body3.paragraphs[0]
        else:
            p = tf_body3.add_paragraph()
        p.text = insight
        p.space_after = Pt(10)
        for r in p.runs:
            r.font.size = Pt(16)
            r.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # --- Slide 4: Summary Slide ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    txBox4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Summary & Next Steps"
    p4.alignment = PP_ALIGN.LEFT
    run4 = p4.runs[0]
    run4.font.size = Pt(28)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

    body4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.0))
    tf_body4 = body4.text_frame
    tf_body4.word_wrap = True

    steps = [
        "Continue investing in Cloud Services infrastructure to sustain growth trajectory.",
        "Evaluate Enterprise Software pricing strategy for mid-market expansion.",
        "Investigate Q3 Hardware Solutions dip and implement preventive measures.",
        "Schedule quarterly review meeting for January 15, 2026.",
        "Prepare detailed product line forecasts for Board presentation."
    ]
    for i, step in enumerate(steps):
        if i == 0:
            p = tf_body4.paragraphs[0]
        else:
            p = tf_body4.add_paragraph()
        p.text = step
        p.space_after = Pt(10)
        for r in p.runs:
            r.font.size = Pt(16)
            r.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
