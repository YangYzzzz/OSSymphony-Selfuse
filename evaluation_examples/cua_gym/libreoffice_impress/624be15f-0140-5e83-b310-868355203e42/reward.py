"""
Reward Script: Bring the 'Logo' image to the front on slide 1
Task ID: impress_tm_056
Domain: libreoffice_impress
Scoring:
  Preconditions (gate): Logo exists as PICTURE, all 4 shapes present
  Component 1 (0.7): Logo is the last shape in spTree (frontmost z-order) on slide 1
  Component 2 (0.3): All other shapes (Title, Subtitle, GradientRectangle) are behind Logo
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_056'


def persist_app_state(domain: str):
    """Save any unsaved GUI edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(1.0)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify that the 'Logo' image has been brought to the front (highest z-order)
    on slide 1 of the presentation.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    if len(prs.slides) < 1:
        print("CRITICAL: No slides in presentation")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]
    shapes = list(slide.shapes)

    # Build a map of shape names for convenience
    shape_names = [s.name for s in shapes]
    print(f"INFO: Slide 1 has {len(shapes)} shapes: {shape_names}")

    # Find the Logo shape
    logo_shape = None
    logo_index = -1
    for i, shape in enumerate(shapes):
        if shape.name == 'Logo':
            logo_shape = shape
            logo_index = i
            break

    # Precondition gate: Logo shape must exist and be a PICTURE
    if logo_shape is None:
        print("FAIL: No shape named 'Logo' found on slide 1")
        print("REWARD: 0.0")
        return 0.0

    if logo_shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        print(f"FAIL: Logo shape type is {logo_shape.shape_type}, expected PICTURE — file may be corrupted")
        print("REWARD: 0.0")
        return 0.0

    # Precondition gate: All 4 original shapes must still be present
    expected_names = {'Logo', 'GradientRectangle', 'Subtitle', 'Title'}
    actual_names = set(shape_names)
    if not expected_names.issubset(actual_names):
        missing = expected_names - actual_names
        print(f"FAIL: Missing shapes: {missing} — shapes may have been deleted")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Preconditions passed — Logo exists as PICTURE, all 4 shapes present")

    # Component 1: Logo is the LAST shape (frontmost z-order) on slide 1 (0.7 points)
    # In python-pptx, shapes list follows spTree order: last shape = rendered on top.
    # In initial_env, Logo is at index 0 (behind all). In golden, it must be last.
    try:
        last_index = len(shapes) - 1
        if logo_index == last_index:
            print(f"PASS: Component 1 — Logo is the last shape (index {logo_index}/{last_index}), frontmost z-order (0.7 pts)")
            total_score += 0.7
        else:
            print(f"FAIL: Component 1 — Logo is at index {logo_index}, expected last index {last_index}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Logo is in front of ALL other shapes (0.3 points)
    # Verify Logo's z-index is strictly higher than Title, Subtitle, and GradientRectangle.
    # This is a stronger check: Logo must be after every other named shape in the list.
    try:
        shapes_not_behind = [
            shape.name for i, shape in enumerate(shapes)
            if shape.name in ('Title', 'Subtitle', 'GradientRectangle') and i >= logo_index
        ]
        for name in shapes_not_behind:
            print(f"  INFO: Shape '{name}' is NOT behind Logo")
        if len(shapes_not_behind) == 0 and logo_index > 0:
            print(f"PASS: Component 2 — All other shapes are behind Logo (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 — Not all other shapes are behind Logo")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
