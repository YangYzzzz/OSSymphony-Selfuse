"""
Initial Setup: Create a presentation with 4 stacked objects on slide 1,
where the Logo image is behind all other objects.
Task ID: impress_tm_056
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_056'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
LOGO_PATH = f'{WORKDIR}/_logo_temp.png'


def create_logo_image():
    """Create a simple logo image for the presentation."""
    img = Image.new('RGBA', (300, 300), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    # Draw a blue circle
    draw.ellipse([20, 20, 280, 280], fill=(41, 98, 255, 255))
    # Draw "AC" text in center
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 100)
    except Exception:
        font = ImageFont.load_default()
    draw.text((75, 80), "AC", fill=(255, 255, 255, 255), font=font)
    img.save(LOGO_PATH)


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    # Create the logo image first
    create_logo_image()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ========== Slide 1: Title slide with 4 stacked objects ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only layout

    # Remove any default placeholders to start clean
    for ph in list(slide1.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # Z-order is determined by add order: first added = bottom, last added = top
    # We want: Logo (bottom) -> Rectangle -> Subtitle -> Title (top)

    # 1. Logo image (BOTTOM - added first, behind everything)
    logo = slide1.shapes.add_picture(
        LOGO_PATH,
        Inches(5.5), Inches(2.0),
        Inches(2.3), Inches(2.3)
    )
    logo.name = "Logo"

    # 2. Gradient rectangle (above logo)
    rect = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2.0), Inches(1.5),
        Inches(9.3), Inches(4.5)
    )
    rect.name = "GradientRectangle"
    fill = rect.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = RGBColor(0x1B, 0x3A, 0x6B)
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = RGBColor(0x3D, 0x7E, 0xC7)
    fill.gradient_stops[1].position = 1.0
    rect.line.fill.background()

    # 3. Subtitle text box (above rectangle)
    sub_box = slide1.shapes.add_textbox(
        Inches(3.0), Inches(4.2),
        Inches(7.3), Inches(1.0)
    )
    sub_box.name = "Subtitle"
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Empowering Innovation Through Technology"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(20)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

    # 4. Title text box (TOP - added last, in front of everything)
    title_box = slide1.shapes.add_textbox(
        Inches(3.0), Inches(2.5),
        Inches(7.3), Inches(1.5)
    )
    title_box.name = "Title"
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Apex Consulting"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # ========== Slide 2: Team Overview ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only
    for ph in list(slide2.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # Title
    t2 = slide2.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.0), Inches(1.0))
    t2.name = "SlideTitle"
    tf = t2.text_frame
    p = tf.paragraphs[0]
    p.text = "Our Leadership Team"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1B, 0x3A, 0x6B)

    # Team members content
    members = [
        ("Elena Rodriguez", "Chief Executive Officer", "15+ years in strategic consulting"),
        ("David Nakamura", "Chief Technology Officer", "Former VP Engineering at TechScale"),
        ("Priya Sharma", "Director of Operations", "Expertise in process optimization"),
        ("James O'Brien", "Head of Client Relations", "Built partnerships with Fortune 500 firms"),
    ]
    y_start = 1.8
    for i, (name, role, desc) in enumerate(members):
        box = slide2.shapes.add_textbox(
            Inches(1.5), Inches(y_start + i * 1.2),
            Inches(10.0), Inches(1.0)
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        run = p.runs[0]
        run.font.name = "Calibri"
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x2C, 0x2C, 0x2C)

        p2 = tf.add_paragraph()
        p2.text = f"{role} - {desc}"
        run2 = p2.runs[0]
        run2.font.name = "Calibri"
        run2.font.size = Pt(16)
        run2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # ========== Slide 3: Services ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only
    for ph in list(slide3.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # Title
    t3 = slide3.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.0), Inches(1.0))
    t3.name = "SlideTitle"
    tf = t3.text_frame
    p = tf.paragraphs[0]
    p.text = "Our Services"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1B, 0x3A, 0x6B)

    services = [
        ("Digital Transformation", "End-to-end digital strategy and implementation"),
        ("Cloud Migration", "Seamless transition to cloud infrastructure"),
        ("Data Analytics", "Actionable insights from complex datasets"),
        ("Cybersecurity", "Comprehensive threat assessment and protection"),
        ("AI Integration", "Machine learning solutions for business processes"),
    ]
    y_start = 1.8
    for i, (svc, desc) in enumerate(services):
        box = slide3.shapes.add_textbox(
            Inches(1.5), Inches(y_start + i * 1.0),
            Inches(10.0), Inches(0.9)
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = svc
        run = p.runs[0]
        run.font.name = "Calibri"
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x2C, 0x2C, 0x2C)

        p2 = tf.add_paragraph()
        p2.text = desc
        run2 = p2.runs[0]
        run2.font.name = "Calibri"
        run2.font.size = Pt(16)
        run2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Clean up temp logo file
    if os.path.exists(LOGO_PATH):
        os.remove(LOGO_PATH)

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
