"""
Initial Setup: 5-slide product roadmap presentation (pre-task state)
Task ID: osworld_impress_bold_all_title_size_underline_004
Domain: libreoffice_impress

Initial state: All text is regular weight (not bold), titles have mixed sizes (not 36pt),
no underlines anywhere.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_bold_all_title_size_underline_004'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_run_font(run, size_pt, bold=False, italic=False, underline=False, color=None):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.underline = underline
    if color:
        run.font.color.rgb = RGBColor(*color)


def create_initial():
    prs = Presentation()
    # Standard widescreen: 13.33 x 7.5 inches
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_layouts = prs.slide_layouts
    # layout 0 = Title Slide, layout 1 = Title and Content, layout 6 = Title Only

    # ---- Slide 1: Title slide ----
    slide1 = prs.slides.add_slide(slide_layouts[0])
    title1 = slide1.shapes.title
    title1.text = ""
    run = title1.text_frame.paragraphs[0].add_run()
    run.text = "NovaTech 2025 Product Roadmap"
    set_run_font(run, size_pt=40, bold=False)  # NOT 36pt, NOT bold

    subtitle = slide1.placeholders[1]
    subtitle.text = ""
    sub_run = subtitle.text_frame.paragraphs[0].add_run()
    sub_run.text = "Strategic Vision & Quarterly Milestones"
    set_run_font(sub_run, size_pt=24, bold=False)

    # ---- Slide 2: Q1 Initiatives ----
    slide2 = prs.slides.add_slide(slide_layouts[1])
    title2 = slide2.shapes.title
    title2.text = ""
    run2 = title2.text_frame.paragraphs[0].add_run()
    run2.text = "Q1 Initiatives: Foundation & Launch"
    set_run_font(run2, size_pt=32, bold=False)  # NOT 36pt, NOT bold

    content2 = slide2.placeholders[1]
    content2.text = ""
    tf2 = content2.text_frame
    tf2.word_wrap = True

    items2 = [
        ("Cloud Infrastructure Overhaul", 20),
        ("Migrate 3 legacy services to AWS EKS by Feb 28", 18),
        ("Target: 99.95% uptime SLA with auto-scaling", 18),
        ("Mobile App v3.0 Beta Release", 20),
        ("Dark mode, redesigned onboarding, push notifications", 18),
        ("Beta group: 5,000 selected users from waitlist", 18),
        ("Security Compliance Certification", 20),
        ("Complete SOC 2 Type II audit by March 15", 18),
    ]
    for i, (text, size) in enumerate(items2):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        r = p.add_run()
        r.text = text
        set_run_font(r, size_pt=size, bold=False)

    # ---- Slide 3: Q2 Growth Strategy ----
    slide3 = prs.slides.add_slide(slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = ""
    run3 = title3.text_frame.paragraphs[0].add_run()
    run3.text = "Q2 Growth Strategy: Expand & Optimize"
    set_run_font(run3, size_pt=28, bold=False)  # NOT 36pt, NOT bold

    content3 = slide3.placeholders[1]
    content3.text = ""
    tf3 = content3.text_frame
    tf3.word_wrap = True

    items3 = [
        ("Market Expansion — APAC Region", 20),
        ("Launch in Singapore, Tokyo, and Sydney", 18),
        ("Localized UI in 4 languages (JP, KO, ZH, EN)", 18),
        ("Partner onboarding target: 150 new resellers", 18),
        ("AI-Powered Analytics Dashboard", 20),
        ("Predictive churn scoring for enterprise accounts", 18),
        ("Real-time usage heatmaps and export to CSV/PDF", 18),
        ("Goal: 40% increase in dashboard DAU by June 30", 18),
    ]
    for i, (text, size) in enumerate(items3):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        r = p.add_run()
        r.text = text
        set_run_font(r, size_pt=size, bold=False)

    # ---- Slide 4: Q3 Platform Maturity ----
    slide4 = prs.slides.add_slide(slide_layouts[1])
    title4 = slide4.shapes.title
    title4.text = ""
    run4 = title4.text_frame.paragraphs[0].add_run()
    run4.text = "Q3 Platform Maturity"
    set_run_font(run4, size_pt=34, bold=False)  # NOT 36pt, NOT bold

    content4 = slide4.placeholders[1]
    content4.text = ""
    tf4 = content4.text_frame
    tf4.word_wrap = True

    items4 = [
        ("API v4 General Availability", 20),
        ("Versioned REST & GraphQL endpoints with rate limiting", 18),
        ("Webhook delivery SLA: 99.9%, latency < 200ms", 18),
        ("Developer portal relaunch with interactive docs", 18),
        ("Enterprise SSO & Governance Module", 20),
        ("SAML 2.0 / OIDC integration for Fortune 500 clients", 18),
        ("Role-based access control across all product surfaces", 18),
        ("Pilot with 12 enterprise accounts in beta cohort", 18),
    ]
    for i, (text, size) in enumerate(items4):
        if i == 0:
            p = tf4.paragraphs[0]
        else:
            p = tf4.add_paragraph()
        r = p.add_run()
        r.text = text
        set_run_font(r, size_pt=size, bold=False)

    # ---- Slide 5: Q4 Wrap-Up & 2026 Vision ----
    slide5 = prs.slides.add_slide(slide_layouts[1])
    title5 = slide5.shapes.title
    title5.text = ""
    run5 = title5.text_frame.paragraphs[0].add_run()
    run5.text = "Q4 Wrap-Up & 2026 Vision"
    set_run_font(run5, size_pt=30, bold=False)  # NOT 36pt, NOT bold

    content5 = slide5.placeholders[1]
    content5.text = ""
    tf5 = content5.text_frame
    tf5.word_wrap = True

    items5 = [
        ("Annual Performance Review", 20),
        ("Revenue target: $48M ARR by December 31", 18),
        ("NPS goal: 72+ (up from 65 in 2024)", 18),
        ("Headcount growth: Engineering +22, Sales +15", 18),
        ("2026 Strategic Bets", 20),
        ("Generative AI features embedded across all product lines", 18),
        ("Global data residency options (EU, US, APAC)", 18),
        ("IPO readiness review with Goldman Sachs advisory team", 18),
    ]
    for i, (text, size) in enumerate(items5):
        if i == 0:
            p = tf5.paragraphs[0]
        else:
            p = tf5.add_paragraph()
        r = p.add_run()
        r.text = text
        set_run_font(r, size_pt=size, bold=False)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open the initial file in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
