"""
Reward Script: Bold all text, set all title font sizes to 36pt, underline only slide 1 title
Task ID: osworld_impress_bold_all_title_size_underline_004
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): All text runs across all slides and shapes are bold
  Component 2 (0.3): All 5 title placeholders have font size == 36pt
  Component 3 (0.3): Slide 1 title is underlined; all other shapes/runs are NOT underlined
"""

import os
from pptx import Presentation
from pptx.util import Pt

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_bold_all_title_size_underline_004'


def get_all_runs(slide):
    """Return list of (shape_name, para_idx, run_idx, run) for all runs in a slide."""
    results = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                for run_idx, run in enumerate(para.runs):
                    if (run.text or "").strip():
                        results.append((shape.name, para_idx, run_idx, run))
    return results


def get_title_shape(slide):
    """Return the title placeholder shape for a slide, or None."""
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 0:
            return shape
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have exactly 5 slides
    if len(prs.slides) != 5:
        print(f"CRITICAL: Expected 5 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # -------------------------------------------------------------------------
    # Component 1: All text runs are bold (0.4 points)
    # Initial state: all runs have bold=False
    # Golden state: all runs have bold=True
    # -------------------------------------------------------------------------
    try:
        non_bold_runs = []
        for slide_idx, slide in enumerate(prs.slides):
            for shape_name, para_idx, run_idx, run in get_all_runs(slide):
                bold = run.font.bold
                # Treat None (inherit) as False — task requires explicit bold=True
                is_bold = bold is True
                if not is_bold:
                    non_bold_runs.append(
                        f"Slide {slide_idx+1}/{shape_name}/para{para_idx}/run{run_idx}: bold={bold}"
                    )

        if not non_bold_runs:
            print("PASS: Component 1 — All text runs are bold (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 — {len(non_bold_runs)} non-bold run(s) found:")
            for detail in non_bold_runs[:5]:
                print(f"  {detail}")
            if len(non_bold_runs) > 5:
                print(f"  ... and {len(non_bold_runs) - 5} more")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # -------------------------------------------------------------------------
    # Component 2: All 5 title placeholders have font size == 36pt (0.3 points)
    # Initial state: title sizes are mixed (40, 32, 28, 34, 30 pt)
    # Golden state: all title sizes are 36pt
    # -------------------------------------------------------------------------
    try:
        EXPECTED_TITLE_SIZE_PT = 36.0
        title_issues = []
        for slide_idx, slide in enumerate(prs.slides):
            title_shape = get_title_shape(slide)
            if title_shape is None:
                title_issues.append(f"Slide {slide_idx+1}: No title placeholder found")
                continue
            for para in title_shape.text_frame.paragraphs:
                for run in para.runs:
                    if not (run.text or "").strip():
                        continue
                    size = run.font.size
                    size_pt = size / 12700 if size is not None else None
                    if size_pt != EXPECTED_TITLE_SIZE_PT:
                        title_issues.append(
                            f"Slide {slide_idx+1} title run: expected {EXPECTED_TITLE_SIZE_PT}pt, found {size_pt}pt"
                        )

        if not title_issues:
            print("PASS: Component 2 — All title font sizes are 36pt (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 — {len(title_issues)} title size issue(s):")
            for detail in title_issues:
                print(f"  {detail}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # -------------------------------------------------------------------------
    # Component 3: Slide 1 title is underlined; all other shapes/runs are NOT underlined (0.3 points)
    # Initial state: no underlines anywhere
    # Golden state: slide 1 title underlined; all others not underlined
    # -------------------------------------------------------------------------
    try:
        underline_issues = []

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                is_title = shape.is_placeholder and shape.placeholder_format.idx == 0
                is_slide1_title = (slide_idx == 0 and is_title)

                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    for run_idx, run in enumerate(para.runs):
                        if not (run.text or "").strip():
                            continue
                        underline = run.font.underline
                        # Normalize: None means not underlined (inherit = False)
                        is_underlined = underline is True

                        if is_slide1_title:
                            # Slide 1 title MUST be underlined
                            if not is_underlined:
                                underline_issues.append(
                                    f"Slide 1 title run{run_idx}: expected underline=True, found {underline}"
                                )
                        else:
                            # All other runs must NOT be underlined
                            if is_underlined:
                                underline_issues.append(
                                    f"Slide {slide_idx+1}/{shape.name}/para{para_idx}/run{run_idx}: "
                                    f"unexpected underline=True"
                                )

        if not underline_issues:
            print("PASS: Component 3 — Slide 1 title underlined; all other runs not underlined (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 3 — {len(underline_issues)} underline issue(s):")
            for detail in underline_issues[:5]:
                print(f"  {detail}")
            if len(underline_issues) > 5:
                print(f"  ... and {len(underline_issues) - 5} more")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(round(total_score, 4), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point: test against the canonical artifact path on the VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
