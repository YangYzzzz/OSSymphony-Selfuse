"""
Initial Setup: Create a 20-slide Product Strategy presentation with no custom slide shows.
Task ID: impress_fix_036
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_036'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    """Add a title slide (layout 0)."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points):
    """Add a slide with title and bullet content (layout 1)."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.paragraphs[0].text = point
        else:
            p = tf.add_paragraph()
            p.text = point
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(prs,
        "NovaTech Product Strategy 2026",
        "Driving Innovation Through Technology & Market Leadership"
    )

    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "Q1 revenue grew 23% YoY to $187M",
        "Customer acquisition cost reduced by 14%",
        "Three new product lines launched successfully",
        "Net Promoter Score improved from 62 to 71",
        "Market share expanded to 18.4% in core segments"
    ])

    # Slide 3: Market Landscape
    add_content_slide(prs, "Market Landscape & Competitive Analysis", [
        "Total addressable market projected at $4.2B by 2027",
        "Key competitors: Apex Solutions (22%), DigitalFirst (15%), Streamline Corp (12%)",
        "Emerging trends: AI-driven automation, edge computing, sustainability",
        "Customer demand shifting toward integrated platform solutions",
        "Regulatory landscape favoring data privacy compliance tools"
    ])

    # Slide 4: Technical Architecture Overview
    add_content_slide(prs, "Technical Architecture Overview", [
        "Microservices-based architecture on Kubernetes (v1.28)",
        "Event-driven messaging via Apache Kafka clusters",
        "PostgreSQL 16 with read replicas for high availability",
        "Redis caching layer achieving 99.97% cache hit rate",
        "GraphQL API gateway handling 12K requests/second"
    ])

    # Slide 5: Platform Engineering Roadmap
    add_content_slide(prs, "Platform Engineering Roadmap", [
        "Phase 1 (Q1): Service mesh migration to Istio 1.20",
        "Phase 2 (Q2): Implement distributed tracing with OpenTelemetry",
        "Phase 3 (Q3): Zero-trust security model deployment",
        "Phase 4 (Q4): Multi-region active-active failover",
        "Estimated infrastructure cost savings: $2.3M annually"
    ])

    # Slide 6: AI/ML Pipeline Infrastructure
    add_content_slide(prs, "AI/ML Pipeline Infrastructure", [
        "MLflow-based model registry with 47 production models",
        "Feature store on Apache Spark processing 850M events/day",
        "GPU cluster: 64x NVIDIA A100 for training workloads",
        "Model serving via TensorFlow Serving with <15ms latency",
        "Automated model retraining triggered by data drift detection"
    ])

    # Slide 7: Customer Segmentation Analysis
    add_content_slide(prs, "Customer Segmentation Analysis", [
        "Enterprise (>1000 employees): 34% of revenue, $145K avg contract",
        "Mid-Market (100-999): 41% of revenue, $38K avg contract",
        "SMB (<100): 25% of revenue, $8.5K avg contract",
        "Highest growth in mid-market segment (+31% YoY)",
        "Enterprise renewal rate at 94%, up from 89%"
    ])

    # Slide 8: Data Infrastructure & Analytics
    add_content_slide(prs, "Data Infrastructure & Analytics", [
        "Data lake on AWS S3 storing 2.4 PB of structured/unstructured data",
        "Real-time analytics with Apache Flink processing streams",
        "dbt transformation layer with 1,200+ tested models",
        "Looker dashboards serving 850 daily active business users",
        "Data quality score improved from 87% to 96%"
    ])

    # Slide 9: Financial Performance Overview
    add_content_slide(prs, "Financial Performance Overview", [
        "Annual recurring revenue: $612M (+27% YoY)",
        "Gross margin: 72.3% (up from 68.9%)",
        "Operating expenses: $284M (46.4% of revenue)",
        "Free cash flow: $89M (14.5% margin)",
        "R&D investment: $124M (20.3% of revenue)"
    ])

    # Slide 10: DevOps & CI/CD Pipeline
    add_content_slide(prs, "DevOps & CI/CD Pipeline Metrics", [
        "Deployment frequency: 47 deploys/day across 23 services",
        "Lead time for changes: 2.1 hours (from 4.8 hours)",
        "Change failure rate: 1.8% (industry benchmark: 7-15%)",
        "Mean time to recovery: 12 minutes (from 45 minutes)",
        "Test coverage: 89.4% with 12,500+ automated test cases"
    ])

    # Slide 11: Go-to-Market Strategy
    add_content_slide(prs, "Go-to-Market Strategy 2026", [
        "Product-led growth model for SMB segment",
        "Enterprise sales team expanding to 85 account executives",
        "Strategic partnerships with Salesforce, ServiceNow, Workday",
        "Channel partner program targeting 200 certified partners",
        "Localization for DACH, Japan, and Southeast Asia markets"
    ])

    # Slide 12: Product Portfolio Overview
    add_content_slide(prs, "Product Portfolio Overview", [
        "NovaTech Core: Workflow automation platform (flagship)",
        "NovaTech Analytics: Business intelligence suite",
        "NovaTech Connect: Integration & API management",
        "NovaTech Shield: Compliance & security module",
        "NovaTech Edge: IoT data processing (new launch Q2 2026)"
    ])

    # Slide 13: Customer Success Metrics
    add_content_slide(prs, "Customer Success Metrics", [
        "Average time to value: 14 days (down from 28 days)",
        "Customer health score: 82/100 across portfolio",
        "Support ticket resolution: 4.2 hours average",
        "Self-service adoption: 67% of support interactions",
        "Customer expansion revenue: 31% of total new ARR"
    ])

    # Slide 14: Security & Compliance Framework
    add_content_slide(prs, "Security & Compliance Framework", [
        "SOC 2 Type II certified (renewed March 2026)",
        "ISO 27001:2022 compliant across all data centers",
        "GDPR and CCPA compliance with automated data mapping",
        "Penetration testing: 0 critical findings in last 4 audits",
        "Bug bounty program with 340+ security researchers"
    ])

    # Slide 15: API & Integration Ecosystem
    add_content_slide(prs, "API & Integration Ecosystem", [
        "RESTful API with 245 endpoints, 99.98% uptime SLA",
        "Webhook delivery: 2.1M events/day with <500ms latency",
        "Pre-built connectors: 180+ enterprise applications",
        "Developer portal: 12,400 registered developers",
        "SDK support: Python, JavaScript, Java, Go, .NET"
    ])

    # Slide 16: Brand & Marketing Strategy
    add_content_slide(prs, "Brand & Marketing Strategy", [
        "Brand awareness increased to 42% in target segments",
        "Content marketing: 3.2M monthly blog visits, 15K newsletter subscribers",
        "Annual conference NovaCon 2026: 8,500 registered attendees",
        "Digital ad spend optimization: CAC reduced to $1,240",
        "Thought leadership: 23 industry publications, 6 analyst reports"
    ])

    # Slide 17: Team & Organizational Growth
    add_content_slide(prs, "Team & Organizational Growth", [
        "Total headcount: 1,847 (up from 1,420)",
        "Engineering team: 680 engineers across 52 squads",
        "Employee satisfaction: 4.3/5.0 (eNPS: +45)",
        "Voluntary turnover: 8.2% (industry avg: 13.5%)",
        "Diversity: 43% women in leadership, 38% URG in engineering"
    ])

    # Slide 18: Cloud Infrastructure & Performance
    add_content_slide(prs, "Cloud Infrastructure & Performance", [
        "Multi-cloud deployment: AWS (primary), GCP (secondary), Azure (DR)",
        "Global CDN with 180+ edge locations, <50ms P95 latency",
        "Auto-scaling handling 3x traffic spikes without degradation",
        "Infrastructure-as-code: 100% Terraform-managed resources",
        "Monthly cloud spend: $1.8M (cost per user down 22%)"
    ])

    # Slide 19: Strategic Partnerships & Alliances
    add_content_slide(prs, "Strategic Partnerships & Alliances", [
        "Technology alliance with AWS: Advanced Partner status",
        "Co-selling agreement with Microsoft generating $18M pipeline",
        "OEM partnership with Accenture for managed services",
        "Academic partnerships with MIT and Stanford for R&D",
        "Industry consortium membership: Cloud Native Computing Foundation"
    ])

    # Slide 20: Vision & Next Steps
    add_content_slide(prs, "Vision & Next Steps for 2026-2027", [
        "Target $800M ARR by end of 2027",
        "Launch AI copilot for workflow automation (Q3 2026)",
        "Expand to 3 new geographic markets",
        "Achieve FedRAMP authorization for government sector",
        "Prepare for potential IPO in H2 2027"
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
