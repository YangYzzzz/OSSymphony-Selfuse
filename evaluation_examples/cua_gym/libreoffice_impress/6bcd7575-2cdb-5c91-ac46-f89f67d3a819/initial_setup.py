"""
Initial Setup: Create a 7-slide project status presentation with no transitions.
Task ID: impress_tm_004
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_004'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_body(slide, title_text, body_lines):
    """Helper to set title and body content on a slide."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    # Find body placeholder (index 1 typically)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            for i, line in enumerate(body_lines):
                if i == 0:
                    tf.paragraphs[0].text = line
                else:
                    p = tf.add_paragraph()
                    p.text = line
            break


def create_initial():
    prs = Presentation()

    # Slide 1: Title Slide (layout 0)
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Project Status Update"
    slide1.placeholders[1].text = "Q1 2026 Review — Prepared by Sarah Chen"

    # Slide 2: Executive Summary (layout 1 - Title + Content)
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide2, "Executive Summary", [
        "Overall project health: On Track",
        "Sprint velocity increased 15% over last quarter",
        "Three major milestones completed ahead of schedule",
        "Client satisfaction score: 4.7 / 5.0",
        "Team expanded from 12 to 16 members",
    ])

    # Slide 3: Budget Overview (layout 1)
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide3, "Budget Overview", [
        "Total Budget: $2,450,000",
        "Spent to Date: $1,187,500 (48.5%)",
        "Remaining: $1,262,500",
        "Infrastructure costs reduced by 22% via cloud migration",
        "Contractor expenses under budget by $45,000",
        "Forecast: Expected to finish 3% under budget",
    ])

    # Slide 4: Timeline (layout 1) — NO TRANSITION
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide4, "Timeline", [
        "Phase 1 — Requirements & Design: Jan 6 – Feb 14 (Complete)",
        "Phase 2 — Core Development: Feb 17 – Apr 25 (In Progress)",
        "Phase 3 — Integration Testing: Apr 28 – May 30",
        "Phase 4 — UAT & Bug Fixes: Jun 2 – Jun 27",
        "Phase 5 — Production Deployment: Jul 1 – Jul 11",
        "Go-Live Date: July 14, 2026",
    ])

    # Slide 5: Team Allocation (layout 1)
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide5, "Team Allocation", [
        "Frontend: Marcus Johnson, Priya Patel, Alex Rivera (3)",
        "Backend: Wei Zhang, Dmitri Volkov, Fatima Al-Hassan (3)",
        "QA: Jennifer Park, Carlos Mendoza (2)",
        "DevOps: Tyler Brooks, Aisha Okafor (2)",
        "Product: Sarah Chen, Michael Torres (2)",
        "Design: Emma Larsson, Raj Gupta, Yuki Tanaka, David Kim (4)",
    ])

    # Slide 6: Risk Assessment (layout 1)
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide6, "Risk Assessment", [
        "HIGH: Third-party API deprecation in May — mitigation plan active",
        "MEDIUM: Key developer on leave Jun 9–20 — backup assigned",
        "MEDIUM: Database migration complexity higher than estimated",
        "LOW: UI framework update may introduce breaking changes",
        "LOW: Vendor contract renewal pending for monitoring tools",
    ])

    # Slide 7: Next Steps (layout 1)
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide7, "Next Steps", [
        "Complete Sprint 8 deliverables by April 11",
        "Finalize API migration strategy — owner: Wei Zhang",
        "Schedule client demo for April 18",
        "Begin integration test environment setup",
        "Submit Q2 budget adjustment request by April 25",
        "Conduct mid-project retrospective on April 30",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
