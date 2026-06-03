"""
Initial Setup: Build a 12-slide Annual Review presentation.
Slide 5 has title 'Revenue Growth Year-over-Year' with empty content area (no chart).
Task ID: impress_ps_018
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_018'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.paragraphs[0].text = point
        else:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
    return slide


def add_blank_title_slide(prs, title_text):
    """Add a slide with just a title and empty content area."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    # Clear the content placeholder to leave it empty
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.paragraphs[0].text = ""
    return slide


def add_table_slide(prs, title_text, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Add title text box
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x62)

    num_rows = len(rows) + 1
    num_cols = len(headers)
    tbl_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.3), Inches(9), Inches(0.4 * num_rows)
    )
    table = tbl_shape.table

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x2E, 0x4A, 0x62)

    for r, row_data in enumerate(rows, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = str(val)

    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(prs, "Annual Review 2025",
                    "Meridian Technologies Inc.\nPrepared by the Finance & Strategy Team")

    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "Revenue grew 36% year-over-year, reaching $5.0M in total annual revenue",
        "Customer base expanded to 1,200+ enterprise clients across 14 countries",
        "Launched 3 major product lines: CloudSync Pro, DataVault Enterprise, and AI Insights",
        "Operating margin improved from 18% to 24% through cost optimization initiatives",
        "Employee headcount grew from 145 to 210, with 92% retention rate",
    ])

    # Slide 3: Key Achievements
    add_content_slide(prs, "Key Achievements", [
        "Won 'Best Enterprise SaaS' award at TechConnect Summit 2025",
        "Secured Series C funding of $28M led by Pinnacle Ventures",
        "Achieved SOC 2 Type II and ISO 27001 certifications",
        "Reduced average customer onboarding time from 14 days to 5 days",
        "Established partnerships with Salesforce, AWS, and Microsoft Azure",
    ])

    # Slide 4: Financial Overview (table)
    add_table_slide(prs, "Financial Overview", [
        "Metric", "Q1 2025", "Q2 2025", "Q3 2025", "Q4 2025", "FY 2025"
    ], [
        ["Revenue", "$1.0M", "$1.2M", "$1.3M", "$1.5M", "$5.0M"],
        ["COGS", "$420K", "$480K", "$510K", "$570K", "$1.98M"],
        ["Gross Margin", "58%", "60%", "61%", "62%", "60%"],
        ["OpEx", "$380K", "$400K", "$410K", "$430K", "$1.62M"],
        ["Net Income", "$200K", "$320K", "$380K", "$500K", "$1.4M"],
    ])

    # Slide 5: Revenue Growth Year-over-Year — EMPTY content area, just title
    add_blank_title_slide(prs, "Revenue Growth Year-over-Year")

    # Slide 6: Market Analysis
    add_content_slide(prs, "Market Analysis", [
        "Total addressable market (TAM) estimated at $12B by 2027",
        "Enterprise SaaS sector growing at 22% CAGR globally",
        "Key competitors: Acme Solutions (32% share), Zenith Corp (18% share)",
        "Meridian's market share increased from 4.2% to 5.8% in 2025",
        "Emerging opportunities in healthcare and financial services verticals",
    ])

    # Slide 7: Product Roadmap
    add_content_slide(prs, "Product Roadmap 2026", [
        "Q1: Launch AI-powered analytics dashboard with predictive modeling",
        "Q2: Release mobile-first client portal for on-the-go management",
        "Q3: Introduce multi-tenant architecture for enterprise scalability",
        "Q4: Deploy advanced security suite with zero-trust framework",
        "Ongoing: Performance optimization targeting 99.99% uptime SLA",
    ])

    # Slide 8: Team Growth
    add_content_slide(prs, "Team Growth & Culture", [
        "Engineering team expanded from 62 to 95 members (+53%)",
        "Opened new offices in Austin, TX and Berlin, Germany",
        "Launched mentorship program pairing 40 junior engineers with senior leads",
        "Employee satisfaction score of 4.6/5.0 in annual engagement survey",
        "Introduced flexible remote-hybrid work model with 3 office days per week",
    ])

    # Slide 9: Customer Success
    add_content_slide(prs, "Customer Success Highlights", [
        "Net Promoter Score (NPS) improved from 52 to 68",
        "Average contract value increased 28% to $42,000 annually",
        "Customer churn rate reduced from 8.2% to 5.1%",
        "Top clients: GlobalTech Industries, Pacific Health Systems, Vertex Financial",
        "98.7% customer support satisfaction rating with <2hr average response time",
    ])

    # Slide 10: Strategic Partnerships
    add_content_slide(prs, "Strategic Partnerships", [
        "AWS Advanced Technology Partner status achieved in March 2025",
        "Microsoft Azure co-sell agreement driving $1.2M pipeline",
        "Salesforce AppExchange integration downloaded 3,400+ times",
        "Channel partner network expanded to 28 resellers across EMEA and APAC",
        "Joint marketing campaigns generated 850 qualified leads in H2 2025",
    ])

    # Slide 11: Risks & Challenges
    add_content_slide(prs, "Risks & Mitigation Strategies", [
        "Talent competition: Increased compensation packages and equity grants",
        "Market saturation: Differentiating through AI and vertical-specific solutions",
        "Regulatory compliance: Dedicated compliance team for GDPR, CCPA, HIPAA",
        "Supply chain for hardware: Diversified vendor relationships across 3 regions",
        "Currency fluctuation: Hedging strategy for EUR and GBP-denominated contracts",
    ])

    # Slide 12: Outlook
    add_content_slide(prs, "2026 Outlook & Goals", [
        "Revenue target: $7.5M (+50% YoY growth)",
        "Expand into 5 new international markets including Japan and Brazil",
        "Achieve 1,800 enterprise clients by end of 2026",
        "IPO readiness assessment scheduled for Q3 2026",
        "Sustainability initiative: Carbon-neutral operations by December 2026",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
