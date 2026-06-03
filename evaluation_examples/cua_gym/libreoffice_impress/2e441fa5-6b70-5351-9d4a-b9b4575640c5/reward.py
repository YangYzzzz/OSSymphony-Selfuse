"""
Reward Script: Year-over-year comparison chart on slide 5
Task ID: impress_ps_018
Domain: libreoffice_impress
Scoring:
  - Component 1: Line chart exists on slide 5 (0.25)
  - Component 2: Two series with correct data values (0.30)
  - Component 3: Categories are Q1-Q4 (0.15)
  - Component 4: Data labels on 2025 Revenue series (0.15)
  - Component 5: Legend present on chart (0.15)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_018'

# Expected data
EXPECTED_2024 = [800000.0, 900000.0, 1000000.0, 1100000.0]
EXPECTED_2025 = [1000000.0, 1200000.0, 1300000.0, 1500000.0]
EXPECTED_CATEGORIES = ['Q1', 'Q2', 'Q3', 'Q4']


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.chart import XL_CHART_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # slide 5, 0-indexed

    # Find chart shape on slide 5
    chart_shape = None
    for shape in slide.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Line chart exists on slide 5 (0.25 points)
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            chart_type = chart.chart_type
            # LINE (65) or LINE_MARKERS (65) or similar line types
            # XL_CHART_TYPE.LINE = 4, LINE_MARKERS = 65, LINE_STACKED = 63, etc.
            is_line = chart_type in (
                4,    # LINE
                65,   # LINE_MARKERS
                63,   # LINE_STACKED
                67,   # LINE_MARKERS_STACKED
                64,   # LINE_STACKED_100
                68,   # LINE_MARKERS_STACKED_100
            )
            if is_line:
                print(f"PASS: Component 1 -- Line chart found on slide 5, type={chart_type} (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 -- Chart on slide 5 is not a line chart, type={chart_type}")
        else:
            print("FAIL: Component 1 -- No chart found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    if chart_shape is None:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    chart = chart_shape.chart

    # Component 2: Two series with correct data values (0.30 points)
    try:
        series_count = len(chart.series)
        if series_count >= 2:
            s0_vals = list(chart.series[0].values)
            s1_vals = list(chart.series[1].values)

            def values_match(actual, expected, tolerance=0.01):
                if len(actual) != len(expected):
                    return False
                return all(
                    abs(a - e) / max(abs(e), 1) <= tolerance
                    for a, e in zip(actual, expected)
                )

            # Check which series maps to which expected data
            # Could be in either order
            match_0_2024 = values_match(s0_vals, EXPECTED_2024)
            match_0_2025 = values_match(s0_vals, EXPECTED_2025)
            match_1_2024 = values_match(s1_vals, EXPECTED_2024)
            match_1_2025 = values_match(s1_vals, EXPECTED_2025)

            both_correct = (match_0_2024 and match_1_2025) or (match_0_2025 and match_1_2024)

            if both_correct:
                print(f"PASS: Component 2 -- Both series have correct values (0.30 pts)")
                total_score += 0.30
            else:
                # Partial: at least one series correct
                one_correct = match_0_2024 or match_0_2025 or match_1_2024 or match_1_2025
                if one_correct:
                    print(f"PARTIAL: Component 2 -- Only one series has correct values (0.15 pts)")
                    total_score += 0.15
                else:
                    print(f"FAIL: Component 2 -- Series values don't match. S0={s0_vals}, S1={s1_vals}")
        else:
            print(f"FAIL: Component 2 -- Expected 2 series, found {series_count}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Categories are Q1-Q4 (0.15 points)
    try:
        plot = chart.plots[0]
        categories = list(plot.categories)
        cats_normalized = [str(c).strip().upper() for c in categories]
        expected_normalized = [c.upper() for c in EXPECTED_CATEGORIES]

        if cats_normalized == expected_normalized:
            print(f"PASS: Component 3 -- Categories are Q1-Q4 (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 3 -- Expected {EXPECTED_CATEGORIES}, found {categories}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Data labels on 2025 Revenue series (0.15 points)
    # Must use XML to check per-series data labels
    try:
        ns = {'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'}
        found_2025_labels = False

        with zipfile.ZipFile(file_path, 'r') as zf:
            chart_files = [f for f in zf.namelist() if 'chart' in f.lower() and f.endswith('.xml')]
            for cf in chart_files:
                data = zf.read(cf).decode()
                if 'lineChart' not in data:
                    continue
                root = ET.fromstring(data)
                for ser in root.findall('.//c:ser', ns):
                    # Get series name
                    tx_v = ser.find('.//c:tx//c:v', ns)
                    ser_name = tx_v.text if tx_v is not None else ''

                    # Get series values to identify which is 2025
                    val_elements = ser.findall('.//c:val//c:v', ns)
                    ser_vals = [float(v.text) for v in val_elements]

                    # Identify 2025 series by name or matching values
                    if '2025' in ser_name or ser_vals == EXPECTED_2025:
                        dLbls = ser.find('c:dLbls', ns)
                        if dLbls is not None:
                            show_val = dLbls.find('c:showVal', ns)
                            if show_val is not None and show_val.get('val') == '1':
                                found_2025_labels = len(ser_vals) > 0

        if found_2025_labels:
            print(f"PASS: Component 4 -- Data labels enabled on 2025 Revenue series (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 -- No data labels found on 2025 Revenue series")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Legend present on chart (0.15 points)
    try:
        if chart.has_legend:
            print(f"PASS: Component 5 -- Chart legend is present (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 5 -- Chart has no legend")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
persist_app_state()

if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
