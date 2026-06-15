"""
Initial Setup: Philosophy Ethics presentation with 8 slides, slide 3 has two-column layout
Task ID: impress_stu_009
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_009'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                color=None, alignment=None):
    """Helper to add a text box with formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tf


def add_bullet_text(text_frame, text, level=0, font_size=14, color=None):
    """Add a bullet paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    run = p.runs[0]
    run.font.size = Pt(font_size)
    if color:
        run.font.color.rgb = color
    return p


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== Slide 1: Title Slide =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # Title
    add_textbox(slide1, Inches(1.5), Inches(1.5), Inches(10), Inches(2),
                "Philosophy of Ethics", font_size=40, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C), alignment=PP_ALIGN.CENTER)
    # Subtitle
    add_textbox(slide1, Inches(2), Inches(4), Inches(9), Inches(1.5),
                "A Comparative Study of Major Ethical Frameworks",
                font_size=24, color=RGBColor(0x4A, 0x4A, 0x4A),
                alignment=PP_ALIGN.CENTER)
    # Author
    add_textbox(slide1, Inches(3), Inches(5.5), Inches(7), Inches(1),
                "Dr. Helena Vasquez  |  Department of Philosophy  |  Fall 2025",
                font_size=16, color=RGBColor(0x66, 0x66, 0x66),
                alignment=PP_ALIGN.CENTER)

    # ===== Slide 2: Introduction to Ethical Theories =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide2, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                "Introduction to Ethical Theories", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))
    tf2 = add_textbox(slide2, Inches(0.8), Inches(1.6), Inches(11), Inches(5),
                      "Ethics is a branch of philosophy that addresses questions about morality.",
                      font_size=18, color=RGBColor(0x33, 0x33, 0x33))
    add_bullet_text(tf2, "Normative ethics: What actions are right or wrong?", level=0, font_size=16)
    add_bullet_text(tf2, "Meta-ethics: What is the nature of moral judgments?", level=0, font_size=16)
    add_bullet_text(tf2, "Applied ethics: How do we apply ethical principles to real-world dilemmas?", level=0, font_size=16)
    add_bullet_text(tf2, "Descriptive ethics: How do people actually behave and reason morally?", level=0, font_size=16)
    add_bullet_text(tf2, "Throughout this course, we examine frameworks that have shaped moral discourse for centuries.", level=0, font_size=16)

    # ===== Slide 3: Two-Column Layout - Utilitarianism vs Deontology =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    add_textbox(slide3, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                "Utilitarianism vs. Deontology", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C), alignment=PP_ALIGN.CENTER)

    # Left Column Header
    add_textbox(slide3, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.7),
                "Utilitarianism", font_size=24, bold=True,
                color=RGBColor(0x2E, 0x74, 0xB5))
    # Left Column Bullets
    tf_left = add_textbox(slide3, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5),
                          "Greatest good for the greatest number", font_size=16,
                          color=RGBColor(0x33, 0x33, 0x33))
    add_bullet_text(tf_left, "Founded by Jeremy Bentham (1748-1832)", level=0, font_size=16)
    add_bullet_text(tf_left, "Refined by John Stuart Mill", level=0, font_size=16)
    add_bullet_text(tf_left, "Consequences determine moral worth", level=0, font_size=16)
    add_bullet_text(tf_left, "Pleasure and pain as moral calculus", level=0, font_size=16)
    add_bullet_text(tf_left, "Act vs. Rule Utilitarianism debate", level=0, font_size=16)

    # Right Column Header
    add_textbox(slide3, Inches(7), Inches(1.5), Inches(5.5), Inches(0.7),
                "Deontology", font_size=24, bold=True,
                color=RGBColor(0xC0, 0x39, 0x2B))
    # Right Column Bullets
    tf_right = add_textbox(slide3, Inches(7), Inches(2.3), Inches(5.5), Inches(4.5),
                           "Actions are inherently right or wrong", font_size=16,
                           color=RGBColor(0x33, 0x33, 0x33))
    add_bullet_text(tf_right, "Immanuel Kant's Categorical Imperative", level=0, font_size=16)
    add_bullet_text(tf_right, "Duty-based moral reasoning", level=0, font_size=16)
    add_bullet_text(tf_right, "Universal moral laws apply to all", level=0, font_size=16)
    add_bullet_text(tf_right, "Respect for persons as ends, not means", level=0, font_size=16)
    add_bullet_text(tf_right, "W.D. Ross's prima facie duties", level=0, font_size=16)

    # ===== Slide 4: Virtue Ethics =====
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide4, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                "Virtue Ethics", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))
    tf4 = add_textbox(slide4, Inches(0.8), Inches(1.6), Inches(11), Inches(5),
                      "Character-centered approach to morality, rooted in Aristotelian thought.",
                      font_size=18, color=RGBColor(0x33, 0x33, 0x33))
    add_bullet_text(tf4, "Eudaimonia: flourishing as the highest human good", level=0, font_size=16)
    add_bullet_text(tf4, "The Golden Mean: virtue as a balance between extremes", level=0, font_size=16)
    add_bullet_text(tf4, "Practical wisdom (phronesis) guides moral decisions", level=0, font_size=16)
    add_bullet_text(tf4, "Modern revival by Alasdair MacIntyre and Martha Nussbaum", level=0, font_size=16)
    add_bullet_text(tf4, "Critique: cultural relativism in defining virtues", level=0, font_size=16)

    # ===== Slide 5: Care Ethics =====
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide5, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                "Care Ethics", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))
    tf5 = add_textbox(slide5, Inches(0.8), Inches(1.6), Inches(11), Inches(5),
                      "Emphasizes relationships and responsiveness to others' needs.",
                      font_size=18, color=RGBColor(0x33, 0x33, 0x33))
    add_bullet_text(tf5, "Carol Gilligan's 'In a Different Voice' (1982)", level=0, font_size=16)
    add_bullet_text(tf5, "Nel Noddings: caring as fundamental to ethical life", level=0, font_size=16)
    add_bullet_text(tf5, "Challenges the justice-oriented focus of traditional ethics", level=0, font_size=16)
    add_bullet_text(tf5, "Context-dependent moral reasoning over abstract principles", level=0, font_size=16)
    add_bullet_text(tf5, "Applied to nursing ethics, education policy, social work", level=0, font_size=16)

    # ===== Slide 6: Social Contract Theory =====
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide6, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                "Social Contract Theory", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))
    tf6 = add_textbox(slide6, Inches(0.8), Inches(1.6), Inches(11), Inches(5),
                      "Morality arises from agreements among rational agents in society.",
                      font_size=18, color=RGBColor(0x33, 0x33, 0x33))
    add_bullet_text(tf6, "Thomas Hobbes: life without government is 'nasty, brutish, and short'", level=0, font_size=16)
    add_bullet_text(tf6, "John Locke: natural rights to life, liberty, and property", level=0, font_size=16)
    add_bullet_text(tf6, "Jean-Jacques Rousseau: the general will and popular sovereignty", level=0, font_size=16)
    add_bullet_text(tf6, "John Rawls: the veil of ignorance and justice as fairness", level=0, font_size=16)
    add_bullet_text(tf6, "T.M. Scanlon: contractualism and moral principles", level=0, font_size=16)

    # ===== Slide 7: Applied Ethics Case Studies =====
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide7, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                "Applied Ethics: Case Studies", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))
    tf7 = add_textbox(slide7, Inches(0.8), Inches(1.6), Inches(11), Inches(5),
                      "How do these frameworks apply to real-world dilemmas?",
                      font_size=18, color=RGBColor(0x33, 0x33, 0x33))
    add_bullet_text(tf7, "Trolley Problem: utilitarian vs. deontological responses", level=0, font_size=16)
    add_bullet_text(tf7, "Healthcare rationing during a pandemic", level=0, font_size=16)
    add_bullet_text(tf7, "AI and autonomous vehicle decision-making", level=0, font_size=16)
    add_bullet_text(tf7, "Whistleblowing: loyalty vs. public interest", level=0, font_size=16)
    add_bullet_text(tf7, "Environmental ethics and intergenerational justice", level=0, font_size=16)

    # ===== Slide 8: References & Further Reading =====
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide8, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                "References & Further Reading", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))
    tf8 = add_textbox(slide8, Inches(0.8), Inches(1.6), Inches(11), Inches(5),
                      "Mill, J.S. (1863). Utilitarianism. Parker, Son, and Bourn.",
                      font_size=14, color=RGBColor(0x33, 0x33, 0x33))
    add_bullet_text(tf8, "Kant, I. (1785). Groundwork of the Metaphysics of Morals.", level=0, font_size=14)
    add_bullet_text(tf8, "Aristotle. (c. 340 BCE). Nicomachean Ethics.", level=0, font_size=14)
    add_bullet_text(tf8, "Gilligan, C. (1982). In a Different Voice. Harvard University Press.", level=0, font_size=14)
    add_bullet_text(tf8, "Rawls, J. (1971). A Theory of Justice. Harvard University Press.", level=0, font_size=14)
    add_bullet_text(tf8, "MacIntyre, A. (1981). After Virtue. University of Notre Dame Press.", level=0, font_size=14)
    add_bullet_text(tf8, "Scanlon, T.M. (1998). What We Owe to Each Other. Harvard University Press.", level=0, font_size=14)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
