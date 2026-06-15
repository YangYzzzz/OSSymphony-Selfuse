"""
Reward Script: Duplicate slide 3 and place after slide 6
Task ID: impress_stu_009
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Presentation has 9 slides (was 8)
  Component 2 (0.25): Slide 7 title matches slide 3 title ("Utilitarianism vs. Deontology")
  Component 3 (0.25): Slide 7 has same shape count and two-column text content as slide 3
  Component 4 (0.20): Slides after insertion preserved (slide 8 = "Applied Ethics", slide 9 = "References")
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_009'


def persist_app_state(domain: str):
    """Save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            import time
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    slides = list(prs.slides)

    # Component 1: Presentation has exactly 9 slides (0.30 points)
    # Initial has 8 slides; after duplicating slide 3 and inserting after slide 6, should be 9.
    try:
        if num_slides == 9:
            print(f"PASS: Component 1 — Slide count is 9 (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 1 — Expected 9 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Early exit if slide count is wrong (can't verify position-dependent checks)
    if num_slides < 9:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Slide 7 title matches slide 3 title (0.25 points)
    # Slide 3 has title "Utilitarianism vs. Deontology"; the duplicate at position 7 should match.
    try:
        slide3 = slides[2]
        slide7 = slides[6]

        # Extract title text from both slides (first text shape)
        def get_slide_title(slide):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        return text
            return ""

        title_3 = get_slide_title(slide3)
        title_7 = get_slide_title(slide7)

        if title_3 and title_7 and title_7 == title_3:
            print(f"PASS: Component 2 — Slide 7 title matches slide 3: '{title_7}' (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 — Slide 3 title='{title_3}', Slide 7 title='{title_7}'")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 7 has same shape count and two-column content as slide 3 (0.25 points)
    # Slide 3 has 5 shapes with a two-column layout. The copy should match.
    try:
        slide3 = slides[2]
        slide7 = slides[6]

        shapes3_count = len(slide3.shapes)
        shapes7_count = len(slide7.shapes)

        # Collect all text from each slide for comparison
        def get_all_text(slide):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
            return texts

        texts3 = get_all_text(slide3)
        texts7 = get_all_text(slide7)

        shape_count_match = (shapes7_count == shapes3_count)
        text_match = (texts7 == texts3)

        if shape_count_match and text_match:
            print(f"PASS: Component 3 — Slide 7 has {shapes7_count} shapes matching slide 3 content (0.25 pts)")
            total_score += 0.25
        else:
            if not shape_count_match:
                print(f"FAIL: Component 3 — Shape count mismatch: slide 3 has {shapes3_count}, slide 7 has {shapes7_count}")
            if not text_match:
                print(f"FAIL: Component 3 — Text content mismatch between slide 3 and slide 7")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slides after insertion are preserved (0.20 points)
    # After inserting copy at position 7, original slides 7-8 shift to 8-9.
    # Slide 8 should be "Applied Ethics: Case Studies", slide 9 should be "References & Further Reading"
    try:
        slide8 = slides[7]
        slide9 = slides[8]

        title_8 = ""
        title_9 = ""
        for shape in slide8.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                title_8 = shape.text_frame.text.strip()
                break
        for shape in slide9.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                title_9 = shape.text_frame.text.strip()
                break

        slide8_ok = "Applied Ethics" in title_8
        slide9_ok = "References" in title_9

        if slide8_ok and slide9_ok:
            print(f"PASS: Component 4 — Slides 8-9 preserved: '{title_8[:40]}...', '{title_9[:40]}...' (0.20 pts)")
            total_score += 0.20
        else:
            if not slide8_ok:
                print(f"FAIL: Component 4 — Slide 8 expected 'Applied Ethics...', found '{title_8[:60]}'")
            if not slide9_ok:
                print(f"FAIL: Component 4 — Slide 9 expected 'References...', found '{title_9[:60]}'")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
