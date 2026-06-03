"""
FINAL REWARD SCRIPT - SUCCESS
Task: In LibreOffice Impress, I need slide 138 to automatically play the audio file located at ~/Desktop/intro.mp3 the instant that slide appears—no extra clicks. How can I set that up?
Generated: 2025-09-10 17:03:06
Status: success
Model: azure-o3
Total Steps: 12
"""

import os
import re
import zipfile
import lxml.etree as ET
from pptx import Presentation

def verify_impress_audio_autoplay(file_path: str) -> float:
    """Verify that slide 138 in the given PPTX automatically plays
    ~/Desktop/intro.mp3 as soon as the slide appears (autoplay).

    Scoring (progressive):
      • 0.30 – Slide 138 exists in the file
      • 0.35 – intro.mp3 is referenced on that slide (in the slide XML or its .rels)
      • 0.35 – The media element has autoplay enabled (autoPlay="1" or "true")
    A perfect score of 1.0 is returned only when all three checks pass.
    """

    print(f"Starting verification for: {file_path}")
    total_score = 0.0
    max_score = 1.0

    # ---------- 1. File existence and load ----------
    if not os.path.exists(file_path):
        print("✗ Presentation file not found")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Unable to load presentation: {e}")
        return 0.0

    slide_count = len(prs.slides)
    print(f"Slide count detected: {slide_count}")

    # ---------- 2. Slide 138 existence ----------
    if slide_count >= 138:
        total_score += 0.30
        target_slide = prs.slides[137]          # zero-based index
        print("✓ Slide 138 exists (0.30)")
    else:
        print("✗ Slide 138 is missing – cannot continue checks")
        return round(total_score, 6)

    # Resolve internal slide XML path, e.g. 'ppt/slides/slide138.xml'
    slide_xml_path = target_slide.part.partname.lstrip('/')
    print(f"Resolved slide XML path: {slide_xml_path}")

    # ---------- 3. Read slide XML & relationships ----------
    try:
        with zipfile.ZipFile(file_path) as z:
            slide_xml_bytes = z.read(slide_xml_path)
            slide_xml_text = slide_xml_bytes.decode('utf-8', errors='ignore')

            rel_path = os.path.join(os.path.dirname(slide_xml_path),
                                    '_rels',
                                    os.path.basename(slide_xml_path) + '.rels')
            rel_xml_text = ''
            if rel_path in z.namelist():
                rel_xml_text = z.read(rel_path).decode('utf-8', errors='ignore')
    except Exception as e:
        print(f"✗ Error reading slide XML data: {e}")
        return round(total_score, 6)

    # ---------- 4. Verify intro.mp3 reference ----------
    intro_ref_found = bool(re.search(r'intro\.mp3', slide_xml_text, re.IGNORECASE) or
                            re.search(r'intro\.mp3', rel_xml_text,  re.IGNORECASE))

    if intro_ref_found:
        total_score += 0.35
        print("✓ intro.mp3 is referenced on slide 138 (0.35)")
    else:
        print("✗ intro.mp3 not referenced on slide 138")

    # ---------- 5. Verify autoplay enabled ----------
    autoplay_found = False

    # Quick regex check first
    if re.search(r'autoPlay="(?:1|true)"', slide_xml_text, re.IGNORECASE):
        autoplay_found = True
    else:
        # Fallback to XML parse for <a14:media autoPlay="…">
        try:
            root = ET.fromstring(slide_xml_bytes)
            ns = {'a14': 'http://schemas.microsoft.com/office/drawing/2010/main'}
            for elem in root.findall('.//a14:media', namespaces=ns):
                ap_val = elem.get('autoPlay')
                if ap_val and ap_val.lower() in {'1', 'true'}:
                    autoplay_found = True
                    break
        except Exception as e:
            print(f"  Warning: XML parse issue while checking autoplay: {e}")

    if autoplay_found:
        total_score += 0.35
        print("✓ Autoplay attribute detected on media (0.35)")
    else:
        print("✗ Autoplay attribute NOT detected on media")

    # ---------- 6. Final score ----------
    final_score = round(min(total_score, max_score), 6)  # avoid float imprecision
    print(f"Total accumulated score: {final_score}/{max_score}")
    return final_score

# ----------------- EXECUTION -----------------
if __name__ == "__main__":
    test_file = "/home/user/in_libreoffice_impress_i_need_slide_138_to_automatically_play_the_audio_file_located_at_desktopintro_golden.pptx"
    reward = verify_impress_audio_autoplay(test_file)
    print(f"REWARD: {reward}")
