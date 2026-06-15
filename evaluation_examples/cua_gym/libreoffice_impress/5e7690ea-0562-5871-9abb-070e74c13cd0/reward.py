"""
Reward Script: Vertically center all text in every cell of the table on slide 2.
Task ID: impress_tct_026
Domain: libreoffice_impress
Scoring:
  Component 1 (0.6): Proportion of table cells with vertical center alignment
  Component 2 (0.4): Strict completeness - ALL cells vertically centered
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_026'


def persist_app_state(domain):
    """Save any unsaved GUI edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Cannot import pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 2")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: find the table on slide 2 (index 1)
    slide = prs.slides[1]
    table_shape = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_shape = shape
            break

    if table_shape is None:
        print("FAIL: No table found on slide 2")
        print("REWARD: 0.0")
        return 0.0

    table = table_shape.table
    num_rows = len(table.rows)
    num_cols = len(table.columns)
    total_cells = num_rows * num_cols
    print(f"INFO: Table found on slide 2: {num_rows} rows x {num_cols} cols ({total_cells} cells)")

    # Check vertical alignment for each cell via XML anchor attribute
    # anchor='ctr' means vertically centered, 't' means top, 'b' means bottom
    NS_A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    centered_count = 0

    for r in range(num_rows):
        for c in range(num_cols):
            try:
                cell = table.cell(r, c)
                tc = cell._tc
                tcPr = tc.find(f'{NS_A}tcPr')
                anchor = tcPr.get('anchor') if tcPr is not None else None
                if anchor == 'ctr':
                    centered_count += 1
                else:
                    print(f"  DETAIL: Cell({r},{c}) anchor='{anchor}' (not centered)")
            except Exception as e:
                print(f"  ERROR: Cell({r},{c}): {e}")

    print(f"INFO: {centered_count}/{total_cells} cells vertically centered")

    # Component 1: Proportion of cells with vertical center alignment (0.6 points)
    # Progressive partial credit based on fraction of cells correctly centered
    try:
        if total_cells > 0:
            fraction = centered_count / total_cells
            comp1_score = round(0.6 * fraction, 4)
            if fraction > 0:
                print(f"PASS: Component 1 — {centered_count}/{total_cells} cells centered, fraction={fraction:.2f} ({comp1_score} pts)")
                total_score += comp1_score
            else:
                print(f"FAIL: Component 1 — no cells are vertically centered")
        else:
            print("FAIL: Component 1 — table has 0 cells")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Strict completeness - ALL cells vertically centered (0.4 points)
    # Only awarded if every single cell has anchor='ctr'
    try:
        if centered_count == total_cells and total_cells > 0:
            print(f"PASS: Component 2 — all {total_cells} cells vertically centered ({0.4} pts)")
            total_score += 0.4
        else:
            missing = total_cells - centered_count
            print(f"FAIL: Component 2 — {missing} cells not vertically centered")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist app state before verification
persist_app_state("libreoffice_impress")

# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
