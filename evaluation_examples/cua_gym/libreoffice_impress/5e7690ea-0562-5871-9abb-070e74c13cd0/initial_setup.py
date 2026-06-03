"""
Initial Setup: Create an Agenda presentation with a table on slide 2.
Task ID: impress_tct_026
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_026'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Quarterly Business Review"
    slide1.placeholders[1].text = "Q1 2025 - Strategic Planning & Performance"

    # --- Slide 2: Table Slide (2 columns x 8 rows) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Add a title text box
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Meeting Agenda"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x7A)

    # Table: 8 rows x 2 columns
    rows, cols = 8, 2
    tbl_shape = slide2.shapes.add_table(
        rows, cols,
        Inches(0.8), Inches(1.2),
        Inches(11.0), Inches(5.5)
    )
    table = tbl_shape.table

    # Set column widths
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(7.5)

    # Table data
    data = [
        ["Time Slot", "Agenda Item"],
        ["9:00 - 9:15", "Welcome & Opening Remarks by CEO"],
        ["9:15 - 9:45", "Financial Performance Overview\nRevenue, margins, and cash flow analysis"],
        ["9:45 - 10:15", "Product Roadmap Update"],
        ["10:15 - 10:30", "Coffee Break"],
        ["10:30 - 11:00", "Marketing Campaign Results\nDigital channels and brand awareness metrics"],
        ["11:00 - 11:30", "Operations & Supply Chain Review"],
        ["11:30 - 12:00", "Q&A and Next Steps"],
    ]

    # Varying row heights to make some rows taller
    row_heights = [
        Inches(0.5),   # Header
        Inches(0.55),  # Row 1
        Inches(0.95),  # Row 2 - taller (multi-line)
        Inches(0.55),  # Row 3
        Inches(0.5),   # Row 4
        Inches(0.95),  # Row 5 - taller (multi-line)
        Inches(0.55),  # Row 6
        Inches(0.55),  # Row 7
    ]

    for r_idx in range(rows):
        table.rows[r_idx].height = row_heights[r_idx]
        for c_idx in range(cols):
            cell = table.cell(r_idx, c_idx)
            cell.text = data[r_idx][c_idx]

            # Text is top-aligned by default (MSO_ANCHOR.TOP) - explicitly set to be clear
            cell.vertical_anchor = MSO_ANCHOR.TOP

            # Style the text
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    if r_idx == 0:
                        # Header row styling
                        run.font.bold = True
                        run.font.size = Pt(14)
                        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    else:
                        run.font.size = Pt(12)
                        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

            # Header row background
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x2E, 0x4A, 0x7A)

    # --- Slide 3: Key Topics ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    title_box3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.7))
    tf3 = title_box3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Key Discussion Topics"
    run3 = p3.runs[0]
    run3.font.size = Pt(28)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0x2E, 0x4A, 0x7A)

    content_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(10), Inches(5))
    tf_c = content_box.text_frame
    tf_c.word_wrap = True
    topics = [
        "Revenue growth targets for Q2 and beyond",
        "New product launch timeline and resource allocation",
        "Customer acquisition cost reduction strategies",
        "International expansion opportunities in APAC region",
        "Technology infrastructure modernization plan",
        "Talent retention and recruitment pipeline",
    ]
    for i, topic in enumerate(topics):
        if i == 0:
            p = tf_c.paragraphs[0]
        else:
            p = tf_c.add_paragraph()
        p.text = topic
        p.space_after = Pt(12)
        for run in p.runs:
            run.font.size = Pt(16)

    # --- Slide 4: Closing ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    close_box = slide4.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf4 = close_box.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Thank You"
    p4.alignment = PP_ALIGN.CENTER
    run4 = p4.runs[0]
    run4.font.size = Pt(44)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0x2E, 0x4A, 0x7A)

    p5 = tf4.add_paragraph()
    p5.text = "Next review scheduled for July 15, 2025"
    p5.alignment = PP_ALIGN.CENTER
    for run in p5.runs:
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
