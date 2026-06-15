"""
FINAL REWARD SCRIPT - SUCCESS
Task: I'm sprucing up a presentation and want to make the arrows stand out. Could you guide me on adjusting the line width to exactly 2.0 pt and adding an arrowhead just at the tip?
Generated: 2025-08-07 10:28:32
Status: success
Model: o4-mini
Total Steps: 4
"""

import os
from pptx import Presentation
import xml.etree.ElementTree as ET

def verify_arrows(file_path):
    """
    Verifies that all arrow shapes in the presentation at file_path have:
      1) Line width exactly 2.0 pt
      2) An arrowhead only at the end (headEnd) and none at the beginning (tailEnd)

    Returns a progressive score between 0.0 and 1.0 based on completion:
      - 0.20: File exists
      - 0.20: Presentation loads and arrow shapes detected
      - 0.40: All arrow shapes have line width 2.0 pt (partial credit if some)
      - 0.40: All arrow shapes have arrowhead only at end (partial credit if some)
    """
    total_score = 0.0
    max_score = 1.0
    EMU_PER_PT = 12700  # EMUs per point
    # XML namespaces for parsing
    nsmap = {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
    }

    print(f"Starting arrow verification for: {file_path}")
    # 1. File existence (0.20 points)
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print(f"REWARD: {total_score:.2f}")
        return total_score
    print("✓ File exists (0.20 points)")
    total_score += 0.2

    # 2. Load presentation
    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Error loading PPTX: {e}")
        print(f"REWARD: {total_score:.2f}")
        return total_score

    # 3. Detect arrow shapes (0.20 points)
    arrow_shapes = []  # list of tuples: (shape, headEnd_el, tailEnd_el)
    for slide in prs.slides:
        for shape in slide.shapes:
            # Parse <a:ln> element under shape properties
            ln_el = shape._element.find('.//a:ln', nsmap)
            if ln_el is None:
                continue
            headEnd_el = ln_el.find('a:headEnd', nsmap)
            tailEnd_el = ln_el.find('a:tailEnd', nsmap)
            # Identify shape as arrow if headEnd exists or tailEnd != none
            if headEnd_el is not None or (tailEnd_el is not None and tailEnd_el.get('type') != 'none'):
                arrow_shapes.append((shape, headEnd_el, tailEnd_el))
    if not arrow_shapes:
        print("✗ No arrow shapes with arrowheads found (0.00 points)")
        print(f"REWARD: {total_score:.2f}")
        return total_score
    print(f"✓ Found {len(arrow_shapes)} arrow shapes (0.20 points)")
    total_score += 0.2

    # 4. Verify line width exactly 2.0 pt (0.40 points)
    correct_width_count = 0
    for shape, headEnd_el, tailEnd_el in arrow_shapes:
        width_emu = shape.line.width or 0
        width_pt = width_emu / EMU_PER_PT
        print(f"  Shape ID {shape.shape_id} width: {width_pt:.2f} pt")
        if abs(width_pt - 2.0) < 0.01:
            correct_width_count += 1
    width_score = 0.4 * (correct_width_count / len(arrow_shapes))
    if correct_width_count == len(arrow_shapes):
        print(f"✓ All arrow shapes have width 2.0 pt ({width_score:.2f} points)")
    else:
        print(f"✗ {len(arrow_shapes)-correct_width_count}/{len(arrow_shapes)} shapes not width 2.0 pt ({width_score:.2f} points)")
    total_score += width_score

    # 5. Verify arrowhead only at end (0.40 points)
    correct_ah_count = 0
    for shape, headEnd_el, tailEnd_el in arrow_shapes:
        has_head = headEnd_el is not None and headEnd_el.get('type') != 'none'
        has_tail = tailEnd_el is not None and tailEnd_el.get('type') != 'none'
        print(f"  Shape ID {shape.shape_id} headEnd_type: {headEnd_el.get('type') if headEnd_el is not None else 'None'}, tailEnd_type: {tailEnd_el.get('type') if tailEnd_el is not None else 'None'}")
        if has_head and not has_tail:
            correct_ah_count += 1
    ah_score = 0.4 * (correct_ah_count / len(arrow_shapes))
    if correct_ah_count == len(arrow_shapes):
        print(f"✓ All arrow shapes have arrowhead only at end ({ah_score:.2f} points)")
    else:
        print(f"✗ {len(arrow_shapes)-correct_ah_count}/{len(arrow_shapes)} shapes not properly configured ({ah_score:.2f} points)")
    total_score += ah_score

    # Final score
    final_score = min(total_score, max_score)
    print(f"Total score: {final_score:.2f}/{max_score}")
    print(f"REWARD: {final_score:.2f}")
    return final_score

# Execute verification on the golden PPTX
golden_path = '/home/user/im_sprucing_up_a_presentation_and_want_to_make_the_arrows_stand_out_could_you_guide_me_on_adjusting_.pptx'
verify_arrows(golden_path)
