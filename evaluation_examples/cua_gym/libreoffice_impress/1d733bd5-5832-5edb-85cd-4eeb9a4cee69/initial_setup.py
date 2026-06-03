"""
Initial Setup: Create a 6-slide Scientific Method presentation with slide 4 empty (title only).
Task ID: impress_teach_022
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_022'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a textbox with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    # Standard widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "The Scientific Method"
    slide1.placeholders[1].text = "A Guide for Young Researchers"

    # --- Slide 2: What is Science? ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "What is Science?"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Science is the systematic study of the natural world through observation and experimentation."
    p2 = body2.add_paragraph()
    p2.text = "It relies on evidence-based reasoning and reproducible results to build our understanding of how things work."
    p2.level = 0
    p3 = body2.add_paragraph()
    p3.text = "Key principles include objectivity, skepticism, and peer review."
    p3.level = 0

    # --- Slide 3: Observation & Questions ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Observation & Questions"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Every scientific inquiry begins with careful observation of the natural world."
    p3a = body3.add_paragraph()
    p3a.text = "Observations lead to questions: Why does this happen? What causes this pattern?"
    p3a.level = 0
    p3b = body3.add_paragraph()
    p3b.text = "Good questions are specific, measurable, and testable through experimentation."
    p3b.level = 0
    p3c = body3.add_paragraph()
    p3c.text = "Example: Why do plants grow faster near windows than in closets?"
    p3c.level = 1

    # --- Slide 4: The Scientific Method (TITLE ONLY - no content) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title as a textbox since blank layout has no title placeholder
    add_textbox(slide4, Inches(0.5), Inches(0.3), Inches(12), Inches(1.2),
                "The Scientific Method", font_size=36, bold=True,
                color=RGBColor(0x1F, 0x39, 0x64), alignment=PP_ALIGN.LEFT)

    # --- Slide 5: Experimentation Tips ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Experimentation Tips"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Control variables: Change only one factor at a time to isolate its effect."
    p5a = body5.add_paragraph()
    p5a.text = "Record everything: Keep detailed notes of procedures, measurements, and observations."
    p5a.level = 0
    p5b = body5.add_paragraph()
    p5b.text = "Repeat trials: Run experiments multiple times to ensure consistency of results."
    p5b.level = 0
    p5c = body5.add_paragraph()
    p5c.text = "Use appropriate sample sizes to strengthen statistical significance."
    p5c.level = 0

    # --- Slide 6: Summary ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Summary & Next Steps"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "The scientific method provides a structured framework for investigating questions about the world."
    p6a = body6.add_paragraph()
    p6a.text = "Remember: science is iterative — conclusions often lead to new hypotheses and further experiments."
    p6a.level = 0
    p6b = body6.add_paragraph()
    p6b.text = "Next class: We will design our own experiment using the scientific method steps."
    p6b.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
