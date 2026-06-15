"""
Reward Script: Create process flowchart on slide 4 with labeled rectangles and arrows
Task ID: impress_teach_022
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Three rectangles labeled 'Hypothesis', 'Experiment', 'Conclusion' on slide 4
  Component 2 (0.3): All three rectangles filled with #BBDEFB
  Component 3 (0.3): Two right-pointing arrow shapes between the rectangles
"""

import os
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_022'


def persist_app_state(domain: str):
    """Try to save any unsaved GUI edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_auto_shape_type(shape):
    """Get the preset geometry type of an auto shape."""
    try:
        prstGeom = shape._element.find('.//' + qn('a:prstGeom'))
        if prstGeom is not None:
            return prstGeom.get('prst')
    except Exception:
        pass
    return None


def get_shape_fill_rgb(shape):
    """Get the solid fill RGB color of a shape, or None."""
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def get_shape_text(shape):
    """Get all text from a shape's text frame."""
    if shape.has_text_frame:
        return ''.join(p.text for p in shape.text_frame.paragraphs).strip()
    return ''


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)

    # Collect all auto shapes on slide 4
    rectangles = []  # (shape, text, fill_rgb)
    arrows = []      # (shape, prst_type)

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            prst = get_auto_shape_type(shape)
            text = get_shape_text(shape)
            fill_rgb = get_shape_fill_rgb(shape)

            if prst == 'rect':
                rectangles.append((shape, text, fill_rgb))
            elif prst is not None and 'arrow' in prst.lower():
                arrows.append((shape, prst))

    # Also check for freeform or other arrow-like shapes by name
    # (some implementations use different shape types for arrows)
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            prst = get_auto_shape_type(shape)
            # Already handled above
        elif shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            # Some arrow connectors may be freeforms
            name_lower = shape.name.lower() if shape.name else ''
            if 'arrow' in name_lower:
                arrows.append((shape, 'freeform_arrow'))

    print(f"Found {len(rectangles)} rectangles and {len(arrows)} arrows on slide 4")

    # Component 1: Three rectangles labeled 'Hypothesis', 'Experiment', 'Conclusion' (0.4 points)
    expected_labels = {'Hypothesis', 'Experiment', 'Conclusion'}
    try:
        found_labels = set()
        for shape, text, fill_rgb in rectangles:
            if text in expected_labels:
                found_labels.add(text)
                print(f"  Found rectangle: '{text}'")

        if found_labels == expected_labels:
            print(f"PASS: Component 1 -- All three labeled rectangles found: {found_labels} (0.4 pts)")
            total_score += 0.4
        else:
            missing = expected_labels - found_labels
            print(f"FAIL: Component 1 -- Missing rectangle labels: {missing}")
            # Partial credit: proportional to found labels
            partial = len(found_labels) / len(expected_labels) * 0.4
            if partial > 0:
                print(f"  Partial credit: {partial:.2f} pts for {len(found_labels)}/3 labels")
                total_score += partial
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: All three rectangles filled with #BBDEFB (0.3 points)
    try:
        correctly_filled = 0
        for shape, text, fill_rgb in rectangles:
            if text in expected_labels:
                if fill_rgb and fill_rgb.upper() == 'BBDEFB':
                    correctly_filled += 1
                    print(f"  Fill OK: '{text}' has fill #{fill_rgb}")
                else:
                    print(f"  Fill FAIL: '{text}' has fill {fill_rgb}, expected #BBDEFB")

        if correctly_filled == 3:
            print(f"PASS: Component 2 -- All 3 rectangles have #BBDEFB fill (0.3 pts)")
            total_score += 0.3
        elif correctly_filled > 0:
            partial = correctly_filled / 3 * 0.3
            print(f"PARTIAL: Component 2 -- {correctly_filled}/3 rectangles have correct fill ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 -- No rectangles have #BBDEFB fill")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Two right-pointing arrows between rectangles (0.3 points)
    try:
        # Count arrow-type shapes
        right_arrow_count = len(arrows)
        print(f"  Arrow shapes found: {right_arrow_count}")
        for shape, prst in arrows:
            print(f"    Arrow: name='{shape.name}', type='{prst}'")

        if right_arrow_count >= 2:
            print(f"PASS: Component 3 -- {right_arrow_count} arrow shapes found (0.3 pts)")
            total_score += 0.3
        elif right_arrow_count == 1:
            print(f"PARTIAL: Component 3 -- Only 1 arrow found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 3 -- No arrow shapes found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
