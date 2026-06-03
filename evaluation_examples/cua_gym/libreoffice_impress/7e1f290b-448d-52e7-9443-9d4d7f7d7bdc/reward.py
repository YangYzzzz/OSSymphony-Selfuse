"""
Reward Script: IoT Project Presentation
Task ID: impress_wf_074
Domain: libreoffice_impress
Scoring:
  C1 (0.15) - File on Desktop with 10 slides
  C2 (0.10) - Slide 1 title text 'Smart Home IoT Architecture'
  C3 (0.10) - Dark background #263238 on all slides
  C4 (0.15) - Slide 2: central hub + 6 device shapes with connectors
  C5 (0.10) - Slide 2: Appear animations on device shapes
  C6 (0.10) - Slide 3: hardware components table with specs
  C7 (0.10) - Slide 5: data flow diagram (sensor->gateway->cloud->dashboard)
  C8 (0.05) - Slide 7: nested rectangles for security layers
  C9 (0.10) - Slide 8: bar chart present
  C10 (0.05) - Slide 9: cost breakdown table
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_074'

# Persistence hook for unsaved GUI edits
def persist_app_state():
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        import time
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: File on Desktop with exactly 10 slides (0.15 pts)
    try:
        if num_slides == 10:
            print(f"PASS: Component 1 - File has 10 slides (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 - Expected 10 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Slide 1 has title 'Smart Home IoT Architecture' (0.10 pts)
    try:
        if num_slides >= 1:
            slide1 = prs.slides[0]
            all_text = []
            for shape in slide1.shapes:
                if shape.has_text_frame:
                    all_text.append(shape.text_frame.text.strip())
            found_title = any("smart home iot architecture" in t.lower() for t in all_text)
            if found_title:
                print(f"PASS: Component 2 - Slide 1 title found (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 2 - 'Smart Home IoT Architecture' not found in slide 1 texts: {all_text}")
        else:
            print(f"FAIL: Component 2 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: All slides have dark background #263238 (0.10 pts)
    try:
        if num_slides >= 10:
            dark_count = 0
            for i, slide in enumerate(prs.slides):
                fill = slide.background.fill
                try:
                    if fill.type is not None and str(fill.fore_color.rgb).upper() == "263238":
                        dark_count += 1
                except Exception:
                    pass
            if dark_count >= 8:
                # Allow up to 2 slides to have slightly different bg handling
                print(f"PASS: Component 3 - {dark_count}/10 slides have #263238 background (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 3 - Only {dark_count}/10 slides have #263238 background")
        else:
            print(f"FAIL: Component 3 - Not enough slides for background check")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Slide 2 has central hub circle + 6 device shapes with connectors (0.15 pts)
    try:
        if num_slides >= 2:
            slide2 = prs.slides[1]
            oval_count = 0
            auto_shape_count = 0
            line_count = 0
            device_names = {"thermostat", "camera", "lights", "lock", "speaker", "sensor"}
            found_devices = set()

            for shape in slide2.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    # Check if it's an oval (hub) or rounded rect (device)
                    if "oval" in shape.name.lower():
                        oval_count += 1
                    auto_shape_count += 1
                    if shape.has_text_frame:
                        txt = shape.text_frame.text.strip().lower()
                        for dev in device_names:
                            if dev in txt:
                                found_devices.add(dev)
                elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
                    line_count += 1

            has_hub = oval_count >= 1 or any(
                "hub" in s.text_frame.text.lower()
                for s in slide2.shapes
                if s.has_text_frame and s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            )
            has_devices = len(found_devices) >= 5  # allow 1 miss
            has_connectors = line_count >= 5

            if has_hub and has_devices and has_connectors:
                print(f"PASS: Component 4 - Slide 2 has hub + {len(found_devices)} devices + {line_count} connectors (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 4 - hub={has_hub}, devices={found_devices}, connectors={line_count}")
        else:
            print(f"FAIL: Component 4 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Slide 2 has Appear animations (presetID=1, presetClass=entr) (0.10 pts)
    try:
        ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide2.xml') as f:
                root = ET.parse(f).getroot()
                timing = root.find(f'.//{{{ns_p}}}timing')
                if timing is not None:
                    timing_xml = ET.tostring(timing, encoding='unicode')
                    # Appear animation: presetID="1" presetClass="entr"
                    appear_count = timing_xml.count('presetID="1"')
                    if appear_count >= 4:
                        print(f"PASS: Component 5 - Slide 2 has {appear_count} Appear animations (0.10 pts)")
                        total_score += 0.10
                    else:
                        print(f"FAIL: Component 5 - Only {appear_count} Appear animations found (need >= 4)")
                else:
                    print(f"FAIL: Component 5 - No timing/animation element found on slide 2")
    except KeyError:
        print(f"FAIL: Component 5 - slide2.xml not found in archive")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    # Component 6: Slide 3 has hardware components table with specs (0.10 pts)
    try:
        if num_slides >= 3:
            slide3 = prs.slides[2]
            table_found = 0
            for shape in slide3.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    tbl = shape.table
                    rows = len(tbl.rows)
                    cols = len(tbl.columns)
                    # Should have header row + multiple data rows, multiple columns
                    if rows >= 5 and cols >= 3:
                        # Check header row for expected columns
                        header_texts = [tbl.cell(0, c).text.strip().lower() for c in range(cols)]
                        has_component = any("component" in h for h in header_texts)
                        has_model_or_spec = any(h in ("model", "protocol", "power", "price", "spec", "specs") for h in header_texts)
                        if has_component and has_model_or_spec:
                            table_found += 1
                            print(f"PASS: Component 6 - Slide 3 has {rows}x{cols} hardware table with proper headers (0.10 pts)")
                            total_score += 0.10
            if table_found == 0:
                print(f"FAIL: Component 6 - No suitable hardware components table on slide 3")
        else:
            print(f"FAIL: Component 6 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 6 - {e}")

    # Component 7: Slide 5 has data flow diagram (sensor->gateway->cloud->dashboard) (0.10 pts)
    try:
        if num_slides >= 5:
            slide5 = prs.slides[4]
            flow_keywords = {"sensor", "gateway", "cloud", "dashboard"}
            found_keywords = set()
            line_count = 0
            for shape in slide5.shapes:
                if shape.has_text_frame:
                    txt = shape.text_frame.text.strip().lower()
                    for kw in flow_keywords:
                        if kw in txt:
                            found_keywords.add(kw)
                if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                    line_count += 1

            if len(found_keywords) >= 3 and line_count >= 3:
                print(f"PASS: Component 7 - Slide 5 data flow has {found_keywords} + {line_count} connectors (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 7 - keywords={found_keywords}, connectors={line_count}")
        else:
            print(f"FAIL: Component 7 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 7 - {e}")

    # Component 8: Slide 7 has nested rectangles for security layers (0.05 pts)
    try:
        if num_slides >= 7:
            slide7 = prs.slides[6]
            rect_count = 0
            security_keywords = {"security", "layer", "physical", "network", "application", "data", "encryption"}
            found_sec = set()
            for shape in slide7.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    if "rectangle" in shape.name.lower() or "rect" in shape.name.lower():
                        rect_count += 1
                if shape.has_text_frame:
                    txt = shape.text_frame.text.strip().lower()
                    for kw in security_keywords:
                        if kw in txt:
                            found_sec.add(kw)

            if rect_count >= 3 and len(found_sec) >= 2:
                print(f"PASS: Component 8 - Slide 7 has {rect_count} rectangles + security keywords {found_sec} (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 8 - rects={rect_count}, security keywords={found_sec}")
        else:
            print(f"FAIL: Component 8 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 8 - {e}")

    # Component 9: Slide 8 has a chart (bar chart for power consumption) (0.10 pts)
    try:
        if num_slides >= 8:
            slide8 = prs.slides[7]
            chart_count = 0
            for shape in slide8.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    chart_count += 1
                    break
            if chart_count > 0:
                print(f"PASS: Component 9 - Slide 8 has a chart (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 9 - No chart found on slide 8")
        else:
            print(f"FAIL: Component 9 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 9 - {e}")

    # Component 10: Slide 9 has a cost breakdown table (0.05 pts)
    try:
        if num_slides >= 9:
            slide9 = prs.slides[8]
            cost_table_count = 0
            for shape in slide9.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    tbl = shape.table
                    if len(tbl.rows) >= 3 and len(tbl.columns) >= 2:
                        cost_table_count += 1
                        print(f"PASS: Component 10 - Slide 9 has {len(tbl.rows)}x{len(tbl.columns)} cost table (0.05 pts)")
                        total_score += 0.05
            if cost_table_count == 0:
                print(f"FAIL: Component 10 - No cost table found on slide 9")
        else:
            print(f"FAIL: Component 10 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 10 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entrypoint
persist_app_state()

file_path = f'{WORKDIR}/Desktop/IoT_Architecture.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
