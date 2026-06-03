"""
Reward Script: Set forest green (#228B22) background on all slides and white body text
Task ID: osworld_impress_all_slides_background_012
Domain: libreoffice_impress
Scoring:
  Component 1 (0.50 pts): All 11 slides have solid #228B22 background
  Component 2 (0.30 pts): All body/content text runs have white (#FFFFFF) color
  Component 3 (0.20 pts): Full coverage — slides with non-white text runs = 0
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_all_slides_background_012'

EXPECTED_BG = '228B22'       # Forest green
EXPECTED_TEXT = 'FFFFFF'     # White
TOTAL_SLIDES = 11


def get_slide_background_rgb(slide):
    """
    Return the slide background RGB string (e.g. '228B22'), or None if
    the background is transparent / not a solid fill.
    Falls back to slide-master background for inherited fills.
    """
    fill = slide.background.fill
    if fill.type == 1:  # solid fill directly on slide
        return str(fill.fore_color.rgb)
    elif fill.type == 5:  # inherited from master
        master_fill = slide.slide_layout.slide_master.background.fill
        if master_fill.type == 1:
            return str(master_fill.fore_color.rgb)
    return None


def get_all_text_runs(slide):
    """
    Return all non-empty runs across every text-bearing shape on the slide,
    including nested shapes inside groups.
    """
    def extract_shapes(shape):
        shapes = []
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                shapes.extend(extract_shapes(sub))
        if hasattr(shape, 'text_frame'):
            shapes.append(shape)
        return shapes

    runs = []
    for shape in slide.shapes:
        for s in extract_shapes(shape):
            if s.has_text_frame:
                for para in s.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            runs.append(run)
    return runs


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0

    Scoring breakdown:
      0.50 pts — All 11 slides have #228B22 background
      0.30 pts — All body/content text runs are white (#FFFFFF)
      0.20 pts — No slides contain any non-white text run (completeness)
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    print(f"INFO: Loaded presentation with {num_slides} slides")

    # -----------------------------------------------------------------------
    # Component 1: All 11 slides have forest green (#228B22) background (0.50 pts)
    # -----------------------------------------------------------------------
    try:
        slides_with_green_bg = 0
        non_green_slides = []

        for i, slide in enumerate(prs.slides):
            bg = get_slide_background_rgb(slide)
            if bg == EXPECTED_BG:
                slides_with_green_bg += 1
            else:
                non_green_slides.append((i + 1, bg))

        if slides_with_green_bg == TOTAL_SLIDES:
            print(f"PASS: Component 1 — All {TOTAL_SLIDES} slides have #228B22 background (0.50 pts)")
            total_score += 0.50
        else:
            # Partial: award 0.02 pts per slide (max 0.22 for 11 of them, but cap at 0.50)
            # Use binary: all-or-nothing for simplicity at this threshold
            partial = round(slides_with_green_bg / TOTAL_SLIDES * 0.50, 4)
            print(
                f"FAIL: Component 1 — Only {slides_with_green_bg}/{TOTAL_SLIDES} slides have #228B22 bg "
                f"(partial {partial} pts); missing slides: {non_green_slides}"
            )
            total_score += partial
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # -----------------------------------------------------------------------
    # Component 2: All body/content text runs have white (#FFFFFF) color (0.30 pts)
    # -----------------------------------------------------------------------
    try:
        total_colored_runs = 0
        white_runs = 0
        non_white_examples = []

        for i, slide in enumerate(prs.slides):
            for run in get_all_text_runs(slide):
                if run.font.color.type is not None:
                    total_colored_runs += 1
                    rgb = str(run.font.color.rgb)
                    if rgb == EXPECTED_TEXT:
                        white_runs += 1
                    else:
                        if len(non_white_examples) < 5:
                            non_white_examples.append(
                                f"slide {i+1}: {run.text[:20]!r} => #{rgb}"
                            )

        if total_colored_runs == 0:
            print("FAIL: Component 2 — No explicitly-colored text runs found (cannot confirm white)")
        elif white_runs == total_colored_runs:
            print(
                f"PASS: Component 2 — All {total_colored_runs} colored text runs are white #FFFFFF (0.30 pts)"
            )
            total_score += 0.30
        else:
            partial = round(white_runs / total_colored_runs * 0.30, 4)
            print(
                f"FAIL: Component 2 — {white_runs}/{total_colored_runs} runs are white; "
                f"partial {partial} pts; non-white examples: {non_white_examples}"
            )
            total_score += partial
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # -----------------------------------------------------------------------
    # Component 3: No slides contain any non-white text run (completeness) (0.20 pts)
    # This only passes if EVERY run on EVERY slide is white — a higher bar than
    # Component 2 (which checks ratio), ensuring full coverage slide-by-slide.
    # -----------------------------------------------------------------------
    try:
        slides_with_nonwhite = []

        for i, slide in enumerate(prs.slides):
            for run in get_all_text_runs(slide):
                if run.font.color.type is not None:
                    rgb = str(run.font.color.rgb)
                    if rgb != EXPECTED_TEXT:
                        slides_with_nonwhite.append(i + 1)
                        break  # one offending run per slide is enough

        if len(slides_with_nonwhite) == 0:
            print(f"PASS: Component 3 — No slides contain non-white text runs (0.20 pts)")
            total_score += 0.20
        else:
            print(
                f"FAIL: Component 3 — {len(slides_with_nonwhite)} slides still have non-white text: "
                f"{slides_with_nonwhite}"
            )
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # -----------------------------------------------------------------------
    final_score = min(round(total_score, 4), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in this env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
