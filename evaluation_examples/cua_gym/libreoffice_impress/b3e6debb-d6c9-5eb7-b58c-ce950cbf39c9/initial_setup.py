"""
Initial Setup: Environmental report deck — 11 slides, white backgrounds, dark gray body text
Task ID: osworld_impress_all_slides_background_012
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_all_slides_background_012'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

# Dark gray for body text (NOT white — task requires changing to white)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
TITLE_COLOR = RGBColor(0x1A, 0x1A, 0x1A)
WHITE_BG = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_DARK = RGBColor(0x1F, 0x61, 0x45)   # dark green accent for titles only


def set_slide_bg_white(slide):
    """Set slide background to solid white."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE_BG


def add_title_text(slide, title_text, subtitle_text=None):
    """Add title + optional subtitle text box to a slide."""
    # Title box
    txb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1.2))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.bold = True
    run.font.size = Pt(32)
    run.font.color.rgb = TITLE_COLOR
    run.font.name = "Calibri"

    if subtitle_text:
        txb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(9), Inches(0.8))
        tf2 = txb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle_text
        run2.font.size = Pt(18)
        run2.font.color.rgb = DARK_GRAY
        run2.font.name = "Calibri"


def add_body_text(slide, left, top, width, height, text_lines):
    """Add a body text box with dark gray text."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(14)
        run.font.color.rgb = DARK_GRAY
        run.font.name = "Calibri"


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # -----------------------------------------------------------------------
    # Slide 1 — Cover / Title Slide
    # -----------------------------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    set_slide_bg_white(slide1)

    # Main title
    txb = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.8))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Annual Environmental Impact Report"
    run.font.bold = True
    run.font.size = Pt(40)
    run.font.color.rgb = TITLE_COLOR
    run.font.name = "Calibri"

    # Subtitle
    txb2 = slide1.shapes.add_textbox(Inches(1), Inches(3.4), Inches(8), Inches(1.2))
    tf2 = txb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Greenfield Sustainability Initiative — Fiscal Year 2024"
    run2.font.size = Pt(20)
    run2.font.color.rgb = DARK_GRAY
    run2.font.name = "Calibri"

    txb3 = slide1.shapes.add_textbox(Inches(1), Inches(4.7), Inches(8), Inches(0.7))
    tf3 = txb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "Prepared by the Office of Environmental Affairs"
    run3.font.size = Pt(14)
    run3.font.color.rgb = DARK_GRAY
    run3.font.name = "Calibri"

    # -----------------------------------------------------------------------
    # Slide 2 — Table of Contents
    # -----------------------------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg_white(slide2)
    add_title_text(slide2, "Table of Contents")
    items = [
        "1. Executive Summary",
        "2. Carbon Emissions Overview",
        "3. Renewable Energy Progress",
        "4. Water Conservation Metrics",
        "5. Waste Reduction Initiatives",
        "6. Biodiversity Programs",
        "7. Community Engagement",
        "8. Financial Impact of Green Investments",
        "9. Year-over-Year Comparison",
        "10. Goals and Roadmap for 2025",
    ]
    add_body_text(slide2, Inches(0.8), Inches(1.8), Inches(9), Inches(5), items)

    # -----------------------------------------------------------------------
    # Slide 3 — Executive Summary
    # -----------------------------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg_white(slide3)
    add_title_text(slide3, "Executive Summary")
    body3 = [
        "Greenfield Corp achieved a 14% reduction in total carbon emissions in FY2024.",
        "Renewable energy now accounts for 38% of total operational energy use, up from 29% in FY2023.",
        "Water consumption was reduced by 11% through installation of closed-loop cooling systems.",
        "Over 2,400 tons of waste were diverted from landfills via expanded recycling contracts.",
        "Biodiversity conservation grants totaling $3.2M were distributed across 7 regional programs.",
        "Employee volunteer hours for environmental causes exceeded 18,000 hours — a company record.",
    ]
    add_body_text(slide3, Inches(0.6), Inches(1.9), Inches(9), Inches(4.5), body3)

    # -----------------------------------------------------------------------
    # Slide 4 — Carbon Emissions Overview
    # -----------------------------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg_white(slide4)
    add_title_text(slide4, "Carbon Emissions Overview")
    body4 = [
        "Total Scope 1 Emissions: 42,300 metric tons CO₂e (down from 49,200 in FY2023)",
        "Total Scope 2 Emissions: 18,750 metric tons CO₂e (down from 22,100 in FY2023)",
        "Scope 3 partial estimate: 67,400 metric tons CO₂e (supply chain and business travel)",
        "Key reduction drivers: Fleet electrification (35% EV fleet), on-site solar expansion,",
        "  efficient HVAC upgrades across 14 facilities.",
        "Remaining gap to net-zero target: Approximately 61,050 metric tons CO₂e.",
        "SBTi-aligned pathway: Target net-zero Scope 1+2 by 2035.",
    ]
    add_body_text(slide4, Inches(0.6), Inches(1.9), Inches(9), Inches(4.5), body4)

    # -----------------------------------------------------------------------
    # Slide 5 — Renewable Energy Progress
    # -----------------------------------------------------------------------
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg_white(slide5)
    add_title_text(slide5, "Renewable Energy Progress")
    body5 = [
        "On-site solar capacity: 8.4 MW installed (added 2.1 MW in FY2024)",
        "Wind energy PPAs: 12 MW contracted across three regional utilities",
        "Green tariff agreements: 5 facilities now 100% renewable-powered",
        "Total renewable generation: 48,200 MWh, avoiding ~19,100 metric tons CO₂e",
        "Battery storage pilot: 1.2 MWh system at Northfield facility — 94% uptime",
        "EV charging installations: 320 stations across 18 corporate campuses",
        "FY2025 target: Reach 50% renewable share with two new solar PPA agreements",
    ]
    add_body_text(slide5, Inches(0.6), Inches(1.9), Inches(9), Inches(4.5), body5)

    # -----------------------------------------------------------------------
    # Slide 6 — Water Conservation Metrics
    # -----------------------------------------------------------------------
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg_white(slide6)
    add_title_text(slide6, "Water Conservation Metrics")
    body6 = [
        "Total water withdrawal: 3.14 billion liters (reduced from 3.53 billion in FY2023)",
        "Closed-loop cooling installations at 6 manufacturing plants saved 290M liters",
        "Rainwater harvesting across 9 sites: collected 41M liters for irrigation",
        "Process water recycling rate: 62% (up from 54% in FY2023)",
        "Wastewater compliance: 100% — zero discharge violations across all sites",
        "High water-stress site action plans: Completed for 4 priority locations",
        "FY2025 target: 15% additional reduction via low-flow fixture replacement program",
    ]
    add_body_text(slide6, Inches(0.6), Inches(1.9), Inches(9), Inches(4.5), body6)

    # -----------------------------------------------------------------------
    # Slide 7 — Waste Reduction Initiatives
    # -----------------------------------------------------------------------
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg_white(slide7)
    add_title_text(slide7, "Waste Reduction Initiatives")
    body7 = [
        "Total solid waste generated: 6,870 tons (reduced from 7,340 tons in FY2023)",
        "Landfill diversion rate: 65% — exceeds 60% corporate target",
        "Recycling: 2,850 tons (41%) — metals, cardboard, plastics",
        "Composting & organic diversion: 820 tons (12%) — food waste programs at 11 sites",
        "E-waste: 310 tons properly recycled through certified vendors (R2/e-Stewards)",
        "Hazardous waste: 220 tons — 100% disposed via permitted facilities",
        "Zero-waste-to-landfill certified sites: 3 (Eastport, Maplewood, Riverdale)",
    ]
    add_body_text(slide7, Inches(0.6), Inches(1.9), Inches(9), Inches(4.5), body7)

    # -----------------------------------------------------------------------
    # Slide 8 — Biodiversity Programs
    # -----------------------------------------------------------------------
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg_white(slide8)
    add_title_text(slide8, "Biodiversity Conservation Programs")
    body8 = [
        "Total conservation funding: $3.2M distributed to 7 regional conservation partners",
        "Habitat restoration: 1,450 acres restored or protected across company-owned land",
        "Native plant landscaping: All 18 campuses converted — 42,000 sq ft native plant beds",
        "Wildlife corridor projects: Partnered with 3 land trusts in Pacific Northwest",
        "Pollinator garden installations: 24 gardens across corporate campuses — 8 species observed",
        "Deforestation due-diligence: 100% tier-1 suppliers assessed for forest risk commodities",
        "FY2025 commitment: $4.0M conservation budget, 2,000-acre restoration milestone",
    ]
    add_body_text(slide8, Inches(0.6), Inches(1.9), Inches(9), Inches(4.5), body8)

    # -----------------------------------------------------------------------
    # Slide 9 — Community Engagement
    # -----------------------------------------------------------------------
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg_white(slide9)
    add_title_text(slide9, "Community Engagement")
    body9 = [
        "Employee volunteer hours: 18,240 hours in environmental programs (record high)",
        "Tree-planting drives: 14 events, 9,600 trees planted in collaboration with local municipalities",
        "Environmental education workshops: 38 sessions, reaching 4,200 students in K–12 programs",
        "Community grant program: $1.1M awarded to 22 local environmental nonprofits",
        "Green commute participation: 47% of employees using transit, cycling, or EVs",
        "Internal environmental champions network: 310 certified champions across 14 sites",
        "Annual sustainability report readership: 28,000 downloads — 40% increase YoY",
    ]
    add_body_text(slide9, Inches(0.6), Inches(1.9), Inches(9), Inches(4.5), body9)

    # -----------------------------------------------------------------------
    # Slide 10 — Financial Impact of Green Investments
    # -----------------------------------------------------------------------
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg_white(slide10)
    add_title_text(slide10, "Financial Impact of Green Investments")
    body10 = [
        "Total green capex in FY2024: $47.8M across energy, water, waste, and land programs",
        "Energy efficiency savings: $6.2M/year from HVAC upgrades and LED retrofits",
        "Solar & wind PPA savings: $3.8M vs. grid-rate alternative — payback period ~7.2 years",
        "Water efficiency ROI: $0.9M/year in reduced utility and treatment costs",
        "Avoided carbon credit expenditure: ~$1.4M (at $23/ton avoided rate)",
        "Waste fee avoidance: $520K through recycling and composting programs",
        "Total estimated annual savings: $12.8M — IRR on green investments: ~16%",
    ]
    add_body_text(slide10, Inches(0.6), Inches(1.9), Inches(9), Inches(4.5), body10)

    # -----------------------------------------------------------------------
    # Slide 11 — Goals and Roadmap for 2025
    # -----------------------------------------------------------------------
    slide11 = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg_white(slide11)
    add_title_text(slide11, "Goals and Roadmap for 2025")
    body11 = [
        "Carbon: Achieve 20% Scope 1+2 reduction vs. FY2023 baseline",
        "Energy: Reach 50% renewable share; install 3.0 MW additional solar",
        "Water: Reduce total withdrawal by 15%; expand closed-loop systems to 4 new sites",
        "Waste: Achieve 70% landfill diversion rate; add 2 zero-waste certified sites",
        "Biodiversity: Restore 2,000 acres; launch forest-positive supplier program",
        "Community: 20,000 volunteer hours; expand education program to 60 workshops",
        "Reporting: Publish TCFD-aligned climate risk assessment; ISSB IFRS S2 readiness review",
    ]
    add_body_text(slide11, Inches(0.6), Inches(1.9), Inches(9), Inches(4.5), body11)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(f'libreoffice --impress "{OUTPUT}"'),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
