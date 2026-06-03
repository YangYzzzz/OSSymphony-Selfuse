"""
Reward Script: Interactive Quiz Presentation with Hyperlinks and Color Fills
Task ID: impress_gf4_013
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Shape B (Carbon) on slide 3 has green fill #16A34A
  Component 2 (0.25): Shapes A, C, D on slide 3 have red fill #DC2626
  Component 3 (0.20): Slide 3 answer shapes have correct hyperlinks
                       (B -> slide4, A/C/D -> slide5)
  Component 4 (0.15): Slide 4 has 'Next Question' button linking to slide 6
  Component 5 (0.15): Slide 5 has 'Try Again' button linking to slide 3
"""

import os
from pptx import Presentation
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_013'


def get_shape_fill_color(shape):
    """Extract solid fill color from a shape as uppercase hex string (e.g. 'DC2626')."""
    sp = shape._element
    for child in sp:
        if child.tag.endswith('}spPr'):
            solidFill = child.find(qn('a:solidFill'))
            if solidFill is not None:
                srgb = solidFill.find(qn('a:srgbClr'))
                if srgb is not None:
                    return srgb.get('val', '').upper()
    return None


def get_shape_hyperlink_target(shape, slide):
    """Get the slide filename targeted by the shape's hyperlink action.
    Returns e.g. 'slide4.xml' or None.
    """
    sp = shape._element
    # Look for hlinkClick with ppaction://hlinksldjump
    for hl in sp.iter(qn('a:hlinkClick')):
        action = hl.get('action', '')
        if 'hlinksldjump' in action:
            rid = hl.get(qn('r:id'))
            if rid and rid in slide.part.rels:
                rel = slide.part.rels[rid]
                return str(rel.target_ref)
    return None


def find_shape_by_text_prefix(slide, prefix):
    """Find a shape whose text starts with the given prefix."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text.startswith(prefix):
                return shape
    return None


def find_shape_by_text_exact(slide, text_match):
    """Find a shape whose text matches exactly."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text == text_match:
                return shape
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"CRITICAL: Only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]  # 0-indexed
    slide4 = prs.slides[3]
    slide5 = prs.slides[4]

    # Component 1: Shape B (Carbon) on slide 3 has green fill #16A34A (0.25 points)
    try:
        shape_b = find_shape_by_text_prefix(slide3, 'B:')
        if shape_b is None:
            print("FAIL: Component 1 - Shape 'B: Carbon' not found on slide 3")
        else:
            fill_color = get_shape_fill_color(shape_b)
            if fill_color == '16A34A':
                print(f"PASS: Component 1 - Shape B has green fill #{fill_color} (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 - Shape B fill is #{fill_color}, expected #16A34A")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Shapes A, C, D on slide 3 have red fill #DC2626 (0.25 points)
    try:
        wrong_shapes = []
        wrong_colors_ok = 0
        for prefix in ['A:', 'C:', 'D:']:
            shape = find_shape_by_text_prefix(slide3, prefix)
            if shape is None:
                print(f"FAIL: Component 2 - Shape '{prefix}' not found on slide 3")
                continue
            fill_color = get_shape_fill_color(shape)
            if fill_color == 'DC2626':
                wrong_colors_ok += 1
            else:
                print(f"FAIL: Component 2 - Shape '{prefix}' fill is #{fill_color}, expected #DC2626")
                wrong_shapes.append(prefix)

        if wrong_colors_ok == 3:
            print(f"PASS: Component 2 - All wrong answer shapes (A, C, D) have red fill #DC2626 (0.25 pts)")
            total_score += 0.25
        elif wrong_colors_ok > 0:
            partial = round(0.25 * wrong_colors_ok / 3, 3)
            print(f"PARTIAL: Component 2 - {wrong_colors_ok}/3 wrong shapes have red fill ({partial} pts)")
            total_score += partial
        else:
            print("FAIL: Component 2 - No wrong answer shapes have red fill")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Slide 3 answer shapes have correct hyperlinks (0.20 points)
    # B -> slide4.xml, A/C/D -> slide5.xml
    try:
        links_correct = 0
        total_link_checks = 4

        # Check B links to slide4
        shape_b = find_shape_by_text_prefix(slide3, 'B:')
        if shape_b:
            target = get_shape_hyperlink_target(shape_b, slide3)
            if target and 'slide4' in target:
                links_correct += 1
                print(f"PASS: Component 3 - Shape B links to {target}")
            else:
                print(f"FAIL: Component 3 - Shape B links to {target}, expected slide4.xml")

        # Check A, C, D link to slide5
        for prefix in ['A:', 'C:', 'D:']:
            shape = find_shape_by_text_prefix(slide3, prefix)
            if shape:
                target = get_shape_hyperlink_target(shape, slide3)
                if target and 'slide5' in target:
                    links_correct += 1
                    print(f"PASS: Component 3 - Shape {prefix} links to {target}")
                else:
                    print(f"FAIL: Component 3 - Shape {prefix} links to {target}, expected slide5.xml")

        if links_correct == total_link_checks:
            print(f"PASS: Component 3 - All 4 hyperlinks correct (0.20 pts)")
            total_score += 0.20
        elif links_correct > 0:
            partial = round(0.20 * links_correct / total_link_checks, 3)
            print(f"PARTIAL: Component 3 - {links_correct}/{total_link_checks} links correct ({partial} pts)")
            total_score += partial
        else:
            print("FAIL: Component 3 - No hyperlinks found on slide 3 answer shapes")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Slide 4 has 'Next Question' button linking to slide 6 (0.15 points)
    try:
        next_q_shape = find_shape_by_text_exact(slide4, 'Next Question')
        if next_q_shape is None:
            # Try partial match
            for shape in slide4.shapes:
                if shape.has_text_frame and 'next question' in shape.text_frame.text.strip().lower():
                    next_q_shape = shape
                    break

        if next_q_shape is None:
            print("FAIL: Component 4 - 'Next Question' button not found on slide 4")
        else:
            target = get_shape_hyperlink_target(next_q_shape, slide4)
            if target and 'slide6' in target:
                print(f"PASS: Component 4 - 'Next Question' button links to {target} (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 4 - 'Next Question' button links to {target}, expected slide6.xml")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Slide 5 has 'Try Again' button linking to slide 3 (0.15 points)
    try:
        try_again_shape = find_shape_by_text_exact(slide5, 'Try Again')
        if try_again_shape is None:
            for shape in slide5.shapes:
                if shape.has_text_frame and 'try again' in shape.text_frame.text.strip().lower():
                    try_again_shape = shape
                    break

        if try_again_shape is None:
            print("FAIL: Component 5 - 'Try Again' button not found on slide 5")
        else:
            target = get_shape_hyperlink_target(try_again_shape, slide5)
            if target and 'slide3' in target:
                print(f"PASS: Component 5 - 'Try Again' button links to {target} (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 5 - 'Try Again' button links to {target}, expected slide3.xml")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook: save any unsaved LibreOffice state
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
