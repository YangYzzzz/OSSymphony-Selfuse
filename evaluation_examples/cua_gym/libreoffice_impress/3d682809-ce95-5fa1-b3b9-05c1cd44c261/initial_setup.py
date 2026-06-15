"""
Initial Setup: Interactive quiz presentation with question and answer shapes
Task ID: impress_gf4_013
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

# Install dependencies on VM
subprocess.run(['pip3', 'install', 'python-pptx', 'Pillow'], capture_output=True)

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_013'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_centered_textbox(slide, left, top, width, height, text,
                         font_size=18, bold=False, alignment=PP_ALIGN.CENTER,
                         font_color=None, font_name="Arial"):
    """Helper to add a textbox with centered text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if font_color:
        run.font.color.rgb = font_color
    return txBox


def add_answer_rect(slide, left, top, width, height, text, font_size=16):
    """Add a plain rectangle shape with answer text (no fill color, no hyperlink)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.text = text
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # Light gray fill — neutral, no hint about correctness
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x64, 0x74, 0x8B)  # Slate gray
    return shape


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===================== Slide 1: Title Slide =====================
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)  # Dark navy

    add_centered_textbox(slide1, Inches(2), Inches(2), Inches(9.333), Inches(1.5),
                         "Knowledge Check Quiz", font_size=44, bold=True,
                         font_color=RGBColor(0xFF, 0xFF, 0xFF))
    add_centered_textbox(slide1, Inches(3), Inches(4), Inches(7.333), Inches(1),
                         "Test Your Chemistry Knowledge", font_size=24,
                         font_color=RGBColor(0xBF, 0xDB, 0xFE))

    # ===================== Slide 2: Instructions =====================
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = slide2.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)

    add_centered_textbox(slide2, Inches(2), Inches(0.8), Inches(9.333), Inches(1),
                         "How to Play", font_size=36, bold=True,
                         font_color=RGBColor(0xFF, 0xFF, 0xFF))
    add_centered_textbox(slide2, Inches(2), Inches(2.2), Inches(9.333), Inches(3),
                         "1. Read each question carefully\n"
                         "2. Click on the answer you think is correct\n"
                         "3. If you get it right, you'll move to the next question\n"
                         "4. If you get it wrong, you can try again\n"
                         "5. Good luck!",
                         font_size=20, font_color=RGBColor(0xE2, 0xE8, 0xF0),
                         alignment=PP_ALIGN.LEFT)

    # ===================== Slide 3: Question Slide =====================
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = slide3.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)

    # Question number
    add_centered_textbox(slide3, Inches(1), Inches(0.5), Inches(11.333), Inches(0.6),
                         "Question 1", font_size=18,
                         font_color=RGBColor(0x94, 0xA3, 0xB8))

    # Question text
    add_centered_textbox(slide3, Inches(1.5), Inches(1.2), Inches(10.333), Inches(1.2),
                         "Which element has atomic number 6?",
                         font_size=32, bold=True,
                         font_color=RGBColor(0xFF, 0xFF, 0xFF))

    # Answer shapes in 2x2 grid — NO hyperlinks, NO colored fills (just neutral gray)
    answer_w = Inches(4.5)
    answer_h = Inches(1.5)
    col1_left = Inches(1.8)
    col2_left = Inches(7.0)
    row1_top = Inches(3.2)
    row2_top = Inches(5.2)

    add_answer_rect(slide3, col1_left, row1_top, answer_w, answer_h, "A: Oxygen")
    add_answer_rect(slide3, col2_left, row1_top, answer_w, answer_h, "B: Carbon")
    add_answer_rect(slide3, col1_left, row2_top, answer_w, answer_h, "C: Nitrogen")
    add_answer_rect(slide3, col2_left, row2_top, answer_w, answer_h, "D: Helium")

    # ===================== Slide 4: Correct! =====================
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = slide4.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x05, 0x2E, 0x16)  # Dark green

    add_centered_textbox(slide4, Inches(2), Inches(2.5), Inches(9.333), Inches(1.5),
                         "Correct!", font_size=48, bold=True,
                         font_color=RGBColor(0x4A, 0xDE, 0x80))
    add_centered_textbox(slide4, Inches(3), Inches(4.2), Inches(7.333), Inches(1),
                         "Carbon has atomic number 6. Great job!",
                         font_size=22, font_color=RGBColor(0xBB, 0xF7, 0xD0))
    # No navigation button yet — task requires adding one

    # ===================== Slide 5: Incorrect! =====================
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = slide5.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x45, 0x0A, 0x0A)  # Dark red

    add_centered_textbox(slide5, Inches(2), Inches(2.5), Inches(9.333), Inches(1.5),
                         "Incorrect, try again!", font_size=48, bold=True,
                         font_color=RGBColor(0xFE, 0xCA, 0xCA))
    add_centered_textbox(slide5, Inches(3), Inches(4.2), Inches(7.333), Inches(1),
                         "That's not the right answer. Give it another shot!",
                         font_size=22, font_color=RGBColor(0xFE, 0xE2, 0xE2))
    # No navigation button yet — task requires adding one

    # ===================== Slide 6: Question 2 Placeholder =====================
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = slide6.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)

    add_centered_textbox(slide6, Inches(1), Inches(0.5), Inches(11.333), Inches(0.6),
                         "Question 2", font_size=18,
                         font_color=RGBColor(0x94, 0xA3, 0xB8))
    add_centered_textbox(slide6, Inches(1.5), Inches(1.2), Inches(10.333), Inches(1.2),
                         "What is the chemical symbol for Gold?",
                         font_size=32, bold=True,
                         font_color=RGBColor(0xFF, 0xFF, 0xFF))

    answer_labels = ["A: Go", "B: Gd", "C: Au", "D: Ag"]
    positions = [(col1_left, row1_top), (col2_left, row1_top),
                 (col1_left, row2_top), (col2_left, row2_top)]
    for label, (l, t) in zip(answer_labels, positions):
        add_answer_rect(slide6, l, t, answer_w, answer_h, label)

    # ===================== Slide 7: Fun Facts =====================
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = slide7.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)

    add_centered_textbox(slide7, Inches(2), Inches(0.8), Inches(9.333), Inches(1),
                         "Chemistry Fun Facts", font_size=36, bold=True,
                         font_color=RGBColor(0xFF, 0xFF, 0xFF))
    add_centered_textbox(slide7, Inches(2), Inches(2.2), Inches(9.333), Inches(4),
                         "• The human body contains about 60 chemical elements\n"
                         "• Oxygen is the most abundant element in Earth's crust\n"
                         "• Carbon can form nearly 10 million different compounds\n"
                         "• Fluorine is the most reactive element\n"
                         "• Gallium metal melts in your hand",
                         font_size=20, font_color=RGBColor(0xE2, 0xE8, 0xF0),
                         alignment=PP_ALIGN.LEFT)

    # ===================== Slide 8: End Slide =====================
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = slide8.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)

    add_centered_textbox(slide8, Inches(2), Inches(2.5), Inches(9.333), Inches(1.5),
                         "Quiz Complete!", font_size=44, bold=True,
                         font_color=RGBColor(0xFF, 0xFF, 0xFF))
    add_centered_textbox(slide8, Inches(3), Inches(4.2), Inches(7.333), Inches(1),
                         "Thank you for taking the Knowledge Check Quiz!",
                         font_size=22, font_color=RGBColor(0xBF, 0xDB, 0xFE))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
