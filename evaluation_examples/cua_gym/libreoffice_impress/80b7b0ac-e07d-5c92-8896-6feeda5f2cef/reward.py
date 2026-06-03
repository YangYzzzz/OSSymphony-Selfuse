"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 42 of my LibreOffice Impress presentation, I’ve got an extra placeholder labeled “Text Box 2” beneath the main body text. How do I delete only that secondary box and keep everything else on the slide intact?
Generated: 2025-09-10 12:22:30
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
from pptx import Presentation

def verify_libreoffice_impress_task(file_path):
    """
    Verification script for the task:
    "On slide 42 of my LibreOffice Impress presentation, delete the extra
    placeholder labelled ‘Text Box 2’ beneath the main body text while
    keeping everything else on the slide intact."

    Scoring (progressive, max 1.0):
        • 0.10 – Presentation still has exactly 42 slides (structure intact).
        • 0.50 – Shape named "Text Box 2" (or variant) is GONE from slide 42.
        • 0.20 – Title text on slide 42 is still present and unaltered.
        • 0.20 – Main body text on slide 42 is still present and unaltered.

    Any failure in opening the file or major mismatch yields 0 for that
    criterion. The score is the sum of earned points, capped at 1.0.
    """
    print(f"Verifying presentation: {file_path}")

    total_score = 0.0
    max_score   = 1.0

    # ------------------------------------------------------------------
    # 1. Load the presentation (no points for mere existence / loading)
    # ------------------------------------------------------------------
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Could not open presentation file: {e}")
        return 0.0

    # ------------------------------------------------------------------
    # 2. Basic structural check – slide count remains 42 (0.10 pts)
    # ------------------------------------------------------------------
    slide_count = len(prs.slides)
    print(f"Total slides found: {slide_count}")
    if slide_count == 42:
        total_score += 0.10
        print("✓ Slide count intact (0.10)")
    else:
        print("✗ Slide count changed – expected 42")

    # If slide 42 doesn't exist we cannot continue meaningful checks
    if slide_count < 42:
        print("✗ Slide 42 missing – terminating verification early")
        print(f"Final score: {total_score}")
        return total_score

    slide42 = prs.slides[41]  # zero-based index

    # ------------------------------------------------------------------
    # 3. Verify absence of the extra placeholder ‘Text Box 2’ (0.50 pts)
    # ------------------------------------------------------------------
    shape_names = [getattr(sh, 'name', '') for sh in slide42.shapes]
    print("Shape names on slide 42:", shape_names)

    # Normalise names to catch common variants such as "TextBox 2"
    def normalised(name):
        return name.lower().replace('textbox', 'text box')

    text_box2_present = any('text box 2' in normalised(n) for n in shape_names)
    if not text_box2_present:
        total_score += 0.50
        print("✓ Extra placeholder ‘Text Box 2’ successfully removed (0.50)")
    else:
        print("✗ ‘Text Box 2’ placeholder still present – no points awarded")

    # ------------------------------------------------------------------
    # 4. Confirm essential content is still present (0.20 + 0.20)
    # ------------------------------------------------------------------
    title_intact = False
    body_intact  = False

    for shape in slide42.shapes:
        if hasattr(shape, 'text'):
            text = shape.text.strip().lower()
            if text == 'slide 42 title':
                title_intact = True
            if text == 'main body content on slide 42.':
                body_intact = True

    if title_intact:
        total_score += 0.20
        print("✓ Title text intact (0.20)")
    else:
        print("✗ Title text missing or altered – no points")

    if body_intact:
        total_score += 0.20
        print("✓ Body text intact (0.20)")
    else:
        print("✗ Body text missing or altered – no points")

    # ------------------------------------------------------------------
    # 5. Finalise score (cap at 1.0) and report
    # ------------------------------------------------------------------
    final_score = min(total_score, max_score)
    print(f"Final score: {final_score}")
    return final_score

# ----------------------------------------------------------------------
# ACTUAL EXECUTION – required by evaluation harness
# ----------------------------------------------------------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/on_slide_42_of_my_libreoffice_impress_presentation_ive_got_an_extra_placeholder_labeled_text_box_2_b_golden.pptx"
    reward = verify_libreoffice_impress_task(FILE_PATH)
    print(f"REWARD: {reward}")
