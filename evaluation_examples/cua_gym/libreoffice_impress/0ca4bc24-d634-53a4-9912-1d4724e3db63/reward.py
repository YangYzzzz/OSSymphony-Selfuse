"""
Reward Script: Change presenter notes font to Liberation Sans 13pt
Task ID: impress_ndo_016
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): All non-empty notes runs use 'Liberation Sans' font
  Component 2 (0.5): All non-empty notes runs use 13pt size
"""

import os
from pptx import Presentation
from pptx.util import Pt

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_016'


def persist_app_state(domain: str):
    """Attempt to save any unsaved GUI state."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has 12 slides
    if len(prs.slides) != 12:
        print(f"PRECONDITION FAIL: Expected 12 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Collect all non-empty notes runs across all slides
    all_runs_info = []
    for i, slide in enumerate(prs.slides):
        try:
            ns = slide.notes_slide
            tf = ns.notes_text_frame
            for p in tf.paragraphs:
                for r in p.runs:
                    if (r.text or "").strip():
                        all_runs_info.append({
                            "slide": i + 1,
                            "font_name": r.font.name,
                            "font_size": r.font.size,
                            "text_preview": (r.text or "")[:30]
                        })
        except Exception:
            # Slide has no notes_slide - that's fine, nothing to check
            pass

    if len(all_runs_info) == 0:
        print("FAIL: No non-empty notes runs found in any slide")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Found {len(all_runs_info)} non-empty notes runs across 12 slides")

    # Component 1: All notes runs use 'Liberation Sans' font (0.5 points)
    try:
        correct_font_count = 0
        wrong_font_runs = []
        for info in all_runs_info:
            if info["font_name"] == "Liberation Sans":
                correct_font_count += 1
            else:
                wrong_font_runs.append(
                    f"Slide {info['slide']}: font='{info['font_name']}', text='{info['text_preview']}'"
                )

        font_ratio = correct_font_count / len(all_runs_info)
        if font_ratio == 1.0:
            print(f"PASS: Component 1 -- All {len(all_runs_info)} runs use 'Liberation Sans' (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 1 -- {correct_font_count}/{len(all_runs_info)} runs use 'Liberation Sans' (ratio={font_ratio:.2f})")
            for w in wrong_font_runs[:5]:
                print(f"  Wrong: {w}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: All notes runs use 13pt font size (0.5 points)
    try:
        expected_size = Pt(13)  # 165100 EMU
        correct_size_count = 0
        wrong_size_runs = []
        for info in all_runs_info:
            if info["font_size"] == expected_size:
                correct_size_count += 1
            else:
                actual_pt = info["font_size"] / 12700 if info["font_size"] else None
                wrong_size_runs.append(
                    f"Slide {info['slide']}: size={actual_pt}pt, text='{info['text_preview']}'"
                )

        size_ratio = correct_size_count / len(all_runs_info)
        if size_ratio == 1.0:
            print(f"PASS: Component 2 -- All {len(all_runs_info)} runs use 13pt (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 2 -- {correct_size_count}/{len(all_runs_info)} runs use 13pt (ratio={size_ratio:.2f})")
            for w in wrong_size_runs[:5]:
                print(f"  Wrong: {w}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist app state before verification
persist_app_state("libreoffice_impress")

# Run verification
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
