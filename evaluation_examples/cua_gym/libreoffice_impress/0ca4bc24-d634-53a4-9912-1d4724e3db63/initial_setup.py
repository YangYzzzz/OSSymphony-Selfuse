"""
Initial Setup: Create a 12-slide Annual Review presentation with mixed-font presenter notes.
Task ID: impress_ndo_016
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_016'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_notes(slide, text, font_name, font_size_pt):
    """Add presenter notes to a slide with a specific font."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    # Clear existing paragraphs
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide content definitions ----
    slides_data = [
        {
            "layout": 0,
            "title": "2025 Annual Review",
            "subtitle": "Meridian Technologies Inc.",
            "notes": "Welcome everyone to our annual review. Please hold questions until the end of each section.",
            "notes_font": "Times New Roman",
            "notes_size": 12,
        },
        {
            "layout": 1,
            "title": "Executive Summary",
            "body": "Revenue grew 23% YoY to $847M\nCustomer base expanded to 12,400 accounts\nLaunched 3 new product lines\nOpened offices in Singapore and Berlin",
            "notes": "Emphasize the revenue milestone - this is our strongest growth year since 2019. The Singapore office is already profitable.",
            "notes_font": "Arial",
            "notes_size": 11,
        },
        {
            "layout": 1,
            "title": "Financial Highlights",
            "body": "Q1: $189M (+18%)\nQ2: $204M (+22%)\nQ3: $221M (+26%)\nQ4: $233M (+28%)\nGross Margin: 68.4%\nOperating Income: $142M",
            "notes": "Q4 was particularly strong due to enterprise deal closures. Gross margin improvement driven by automation initiatives in the supply chain.",
            "notes_font": "Times New Roman",
            "notes_size": 14,
        },
        {
            "layout": 1,
            "title": "Product Development",
            "body": "Atlas Platform v3.0 released in March\nNova Analytics suite launched in July\nEdge Computing SDK beta in October\n47 feature releases across all products",
            "notes": "Atlas v3.0 adoption rate exceeded projections by 40%. Nova Analytics already has 2,100 active users. Edge SDK beta feedback has been overwhelmingly positive.",
            "notes_font": "Arial",
            "notes_size": 10,
        },
        {
            "layout": 1,
            "title": "Customer Success Metrics",
            "body": "Net Promoter Score: 72 (up from 64)\nCustomer Retention Rate: 94.2%\nAverage Contract Value: $68,300\nSupport Resolution Time: 4.2 hours",
            "notes": "NPS improvement is significant. Focus on the retention rate - industry average is 85%. Support resolution time decreased by 31% thanks to the new AI-assisted triage system.",
            "notes_font": "Times New Roman",
            "notes_size": 11,
        },
        {
            "layout": 5,
            "title": None,
            "textbox": True,
            "textbox_text": "Engineering Team Achievements\n\n- Migrated 89% of services to Kubernetes\n- Reduced deployment time from 45 min to 8 min\n- Achieved 99.97% uptime across all production systems\n- Implemented zero-trust security architecture",
            "notes": "The Kubernetes migration was led by the platform team under Rachel Torres. Deployment time reduction is a direct result of the new CI/CD pipeline.",
            "notes_font": "Arial",
            "notes_size": 13,
        },
        {
            "layout": 1,
            "title": "Marketing & Brand",
            "body": "Website traffic: 4.8M monthly visitors (+35%)\nSocial media followers: 890K (+52%)\nBrand awareness: 41% in target markets\nMarketing qualified leads: 28,400",
            "notes": "The brand awareness survey was conducted across North America, Europe, and APAC. Social media growth was primarily driven by the thought leadership campaign on LinkedIn.",
            "notes_font": "Times New Roman",
            "notes_size": 12,
        },
        {
            "layout": 1,
            "title": "Talent & Culture",
            "body": "Total headcount: 1,847 (+312 new hires)\nEmployee satisfaction: 4.3/5.0\nVoluntary turnover: 8.7%\nDiversity hiring: 46% of new hires from underrepresented groups",
            "notes": "Turnover rate is well below industry benchmark of 13%. The mentorship program launched in Q2 has been a key driver of employee satisfaction scores.",
            "notes_font": "Arial",
            "notes_size": 14,
        },
        {
            "layout": 1,
            "title": "Strategic Partnerships",
            "body": "Partnership with CloudVista for AI integration\nJoint venture with Hanover Group in EMEA\nTechnology alliance with Nexus Systems\nResearch collaboration with MIT and Stanford",
            "notes": "The CloudVista partnership is expected to generate $35M in co-sell revenue next year. Hanover JV gives us direct access to the German and Nordic markets.",
            "notes_font": "Times New Roman",
            "notes_size": 10,
        },
        {
            "layout": 1,
            "title": "Risk Assessment",
            "body": "Competitive pressure from Apex Solutions (new market entrant)\nRegulatory changes in EU data sovereignty\nSupply chain dependencies on 3 key vendors\nCybersecurity threat landscape evolution",
            "notes": "Apex Solutions raised $200M in Series C and is aggressively hiring our engineers. We need to accelerate our retention initiatives. EU data sovereignty requirements will impact our cloud architecture.",
            "notes_font": "Arial",
            "notes_size": 11,
        },
        {
            "layout": 1,
            "title": "2026 Strategic Priorities",
            "body": "Achieve $1B revenue milestone\nExpand into LATAM and Middle East markets\nLaunch AI-native product suite\nAcquire 2-3 complementary technology companies\nAchieve carbon neutrality in operations",
            "notes": "The $1B target requires 18% growth which is conservative given our trajectory. LATAM expansion will be led by the new regional VP starting in January. AI-native suite budget approved at $45M.",
            "notes_font": "Times New Roman",
            "notes_size": 13,
        },
        {
            "layout": 0,
            "title": "Thank You",
            "subtitle": "Questions & Discussion\ncontact@meridiantech.com",
            "notes": "Open the floor for questions. Remind everyone that detailed department reports are available on the internal wiki. Next all-hands is scheduled for February 15th.",
            "notes_font": "Arial",
            "notes_size": 12,
        },
    ]

    for i, sd in enumerate(slides_data):
        layout_idx = sd["layout"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Set title if present
        if sd.get("title") and slide.shapes.title:
            slide.shapes.title.text = sd["title"]

        # Set subtitle or body
        if sd.get("subtitle") and layout_idx == 0:
            slide.placeholders[1].text = sd["subtitle"]
        elif sd.get("body") and layout_idx == 1:
            slide.placeholders[1].text = sd["body"]

        # Add textbox for blank layout slides
        if sd.get("textbox"):
            txBox = slide.shapes.add_textbox(
                Inches(1), Inches(0.8), Inches(11), Inches(5.5)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = ""
            for line_idx, line in enumerate(sd["textbox_text"].split("\n")):
                if line_idx == 0:
                    tf.paragraphs[0].text = line
                    tf.paragraphs[0].runs[0].font.size = Pt(28)
                    tf.paragraphs[0].runs[0].font.bold = True
                else:
                    p = tf.add_paragraph()
                    p.text = line
                    if line.strip():
                        p.runs[0].font.size = Pt(18)

        # Add presenter notes with specific font
        add_notes(slide, sd["notes"], sd["notes_font"], sd["notes_size"])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
