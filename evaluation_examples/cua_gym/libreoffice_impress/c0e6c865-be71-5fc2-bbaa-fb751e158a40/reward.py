"""
FINAL REWARD SCRIPT - SUCCESS
Task: Reject tracked deletions but keep insertions.
Generated: 2025-10-17 08:09:41
Status: success
Model: azure-o3
Total Steps: 9
"""

import os
import zipfile
from pptx import Presentation

"""
Reward Script: Reject Tracked Deletions but Keep Insertions
----------------------------------------------------------
This script awards a progressive score (0.0 – 1.0) based on three
verifiable requirements:

1.  Inserted text (keyword: "updated") is still present          – 0.4 pts
2.  Deleted text (keyword: "outdated") is also present           – 0.4 pts
    (meaning the deletion was *rejected*)
3.  No residual track-changes markup remains inside PPTX XML      – 0.2 pts

All checks are falsifiable – the score only increases when a specific
requirement is met.  Hard-coded success values are **not** used.
"""

# Keywords to search for (can be adapted per task specifics)
INSERTION_KEYWORD = "updated"
DELETION_KEYWORD  = "outdated"


def _collect_all_slide_texts(prs):
    """Return a list with ALL non-empty text strings from every slide."""
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text.strip()
                if txt:
                    texts.append(txt)
    return texts


def _verify_keyword_present(all_texts, keyword):
    """Case-insensitive keyword search across all slide texts."""
    haystack = " ".join(all_texts).lower()
    return keyword.lower() in haystack


def _verify_no_track_change_markup(file_path):
    """Search PPTX XML parts for common track-changes attributes/tags."""
    tc_patterns = [
        b"o14:deleted", b"o14:ins",        # Office 2010 TC attrs
        b":del", b":ins",                 # Generic TC tags
        b"w14:insert", b"w14:delete",      # Wordprocessing ML variants
        b"v14"                              # Older TC ns prefix
    ]

    with zipfile.ZipFile(file_path, "r") as zf:
        for name in zf.namelist():
            if not (name.startswith("ppt/") and name.endswith(".xml")):
                continue  # Only scan slide-related XML
            data = zf.read(name)
            if any(pat in data for pat in tc_patterns):
                print(f"    ✗ Track-change markup still present in {name}")
                return False
    return True


def calculate_reward(file_path):
    """Main verification routine – returns a float score 0.0–1.0."""
    total, max_score = 0.0, 1.0

    if not os.path.isfile(file_path):
        print("✗ Presentation file not found.")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Unable to load PPTX: {e}")
        return 0.0

    all_texts = _collect_all_slide_texts(prs)
    if not all_texts:
        print("✗ No slide text found – cannot verify task.")
        return 0.0

    # 1) Insertions kept (0.4)
    if _verify_keyword_present(all_texts, INSERTION_KEYWORD):
        print("✓ Inserted content retained (")
        total += 0.4
    else:
        print("✗ Inserted content missing – insertion may have been rejected.")

    # 2) Deletions rejected (deleted word still visible) (0.4)
    if _verify_keyword_present(all_texts, DELETION_KEYWORD):
        print("✓ Deleted content retained as expected (deletion rejected).")
        total += 0.4
    else:
        print("✗ Deleted content not present – deletion may have been accepted.")

    # 3) No residual track-change markup (0.2)
    if _verify_no_track_change_markup(file_path):
        print("✓ No residual track-changes markup found in XML parts.")
        total += 0.2
    else:
        print("✗ Residual track-changes markup detected.")

    final_score = min(total, max_score)
    print(f"Total score: {final_score}")
    return final_score


if __name__ == "__main__":
    FILE = "/home/user/reject_tracked_deletions_but_keep_insertions.pptx"
    reward = calculate_reward(FILE)
    print(f"REWARD: {reward}")
