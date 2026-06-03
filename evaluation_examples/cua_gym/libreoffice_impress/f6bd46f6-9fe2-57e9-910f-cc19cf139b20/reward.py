"""
Reward Script: Apply light blue header and white data row fills to table on slide 2
Task ID: impress_tct_004
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): Header row (row 0) has #D6EAF8 solid fill on all 5 cells
  Component 2 (0.5): Data rows (rows 1-7) have #FFFFFF solid fill on all 35 cells
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_004'


def get_cell_fill_color(cell):
    """
    Extract the solid fill srgbClr hex value from a table cell.
    Returns uppercase hex string (e.g. 'D6EAF8') or None if no solid fill.
    """
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        return None
    solidFill = tcPr.find(qn('a:solidFill'))
    if solidFill is None:
        return None
    srgb = solidFill.find(qn('a:srgbClr'))
    if srgb is not None:
        return srgb.get('val', '').upper()
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Ensure we have at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Expected at least 2 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[1]  # slide 2 (0-indexed)

    # Find the table on slide 2
    table = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            break

    if table is None:
        print("FAIL: No table found on slide 2")
        print("REWARD: 0.0")
        return 0.0

    num_rows = len(table.rows)
    num_cols = len(table.columns)
    print(f"INFO: Table found with {num_rows} rows x {num_cols} cols")

    if num_rows < 2:
        print("FAIL: Table has fewer than 2 rows, cannot verify header vs data")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Header row (row 0) has #D6EAF8 solid fill (0.5 points)
    try:
        header_pass = 0
        for c in range(num_cols):
            cell = table.cell(0, c)
            color = get_cell_fill_color(cell)
            if color == 'D6EAF8':
                header_pass += 1
            else:
                print(f"FAIL: Header cell (0, {c}) expected #D6EAF8, found {color}")

        if header_pass == num_cols:
            print(f"PASS: Component 1 -- All {num_cols} header cells have #D6EAF8 fill (0.5 pts)")
            total_score += 0.5
        elif header_pass > 0:
            partial = 0.5 * (header_pass / num_cols)
            print(f"PARTIAL: Component 1 -- {header_pass}/{num_cols} header cells correct ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- No header cells have #D6EAF8 fill")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Data rows (rows 1 through end) have #FFFFFF solid fill (0.5 points)
    try:
        data_pass = 0
        data_total = (num_rows - 1) * num_cols
        for r in range(1, num_rows):
            for c in range(num_cols):
                cell = table.cell(r, c)
                color = get_cell_fill_color(cell)
                if color == 'FFFFFF':
                    data_pass += 1
                else:
                    print(f"FAIL: Data cell ({r}, {c}) expected #FFFFFF, found {color}")

        if data_pass == data_total:
            print(f"PASS: Component 2 -- All {data_total} data cells have #FFFFFF fill (0.5 pts)")
            total_score += 0.5
        elif data_pass > 0:
            partial = 0.5 * (data_pass / data_total)
            print(f"PARTIAL: Component 2 -- {data_pass}/{data_total} data cells correct ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 -- No data cells have #FFFFFF fill")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
