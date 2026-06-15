"""
Initial Setup: Create a 4-slide presentation with an unformatted table on slide 2.
Task ID: impress_tct_004
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_004'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_cell_text(cell, text, font_name="Calibri", font_size=Pt(14),
                  bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to set cell text with formatting."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Team Roster"
    slide1.placeholders[1].text = "Engineering Division - Q2 2025"

    # --- Slide 2: Table Slide (5 columns x 8 rows, no cell fills) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Add a title textbox
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Engineering Team Members"
    run.font.name = "Calibri"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # Table: 8 rows x 5 columns
    rows, cols = 8, 5
    tbl_shape = slide2.shapes.add_table(
        rows, cols,
        Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.0)
    )
    table = tbl_shape.table

    # Set column widths
    table.columns[0].width = Inches(2.5)  # Name
    table.columns[1].width = Inches(2.0)  # Role
    table.columns[2].width = Inches(2.5)  # Email
    table.columns[3].width = Inches(2.3)  # Start Date
    table.columns[4].width = Inches(3.0)  # Location

    # Headers
    headers = ["Name", "Role", "Email", "Start Date", "Location"]
    for c, h in enumerate(headers):
        set_cell_text(table.cell(0, c), h, bold=True, font_size=Pt(14),
                      color=RGBColor(0x2C, 0x3E, 0x50))

    # Data rows (realistic content, NO background fills)
    data = [
        ["Sarah Chen", "Senior Engineer", "s.chen@example.com", "2023-01-15", "San Francisco, CA"],
        ["Marcus Johnson", "Tech Lead", "m.johnson@example.com", "2022-06-01", "New York, NY"],
        ["Priya Patel", "Backend Developer", "p.patel@example.com", "2024-03-10", "Austin, TX"],
        ["James O'Brien", "DevOps Engineer", "j.obrien@example.com", "2023-08-22", "Seattle, WA"],
        ["Mei Lin", "Frontend Developer", "m.lin@example.com", "2024-01-08", "Chicago, IL"],
        ["Carlos Rivera", "QA Engineer", "c.rivera@example.com", "2023-11-05", "Denver, CO"],
        ["Aisha Mohamud", "Data Engineer", "a.mohamud@example.com", "2024-06-17", "Boston, MA"],
    ]
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            set_cell_text(table.cell(r, c), val, font_size=Pt(12))

    # --- Slide 3: Team Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.7))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = "Team Overview"
    run3.font.name = "Calibri"
    run3.font.size = Pt(28)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    txBox3b = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(11), Inches(5))
    tf3b = txBox3b.text_frame
    tf3b.word_wrap = True
    bullet_items = [
        "Total team size: 7 engineers across 5 US cities",
        "Average tenure: 1.8 years",
        "Primary focus: Cloud infrastructure and platform services",
        "Key initiatives: Migration to Kubernetes, CI/CD pipeline modernization",
        "Hiring targets: 3 additional engineers by end of Q3 2025",
    ]
    for i, item in enumerate(bullet_items):
        if i == 0:
            p = tf3b.paragraphs[0]
        else:
            p = tf3b.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.name = "Calibri"
        run.font.size = Pt(16)
        p.level = 0

    # --- Slide 4: Contact Information ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.7))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    run4 = p4.add_run()
    run4.text = "Contact & Resources"
    run4.font.name = "Calibri"
    run4.font.size = Pt(28)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    txBox4b = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(11), Inches(4))
    tf4b = txBox4b.text_frame
    tf4b.word_wrap = True
    info_items = [
        "Team Slack channel: #eng-platform",
        "Weekly standup: Tuesday & Thursday 10:00 AM PT",
        "Sprint planning: Every other Monday 2:00 PM PT",
        "Documentation wiki: https://wiki.internal/eng-platform",
        "On-call rotation: PagerDuty schedule 'Platform-Primary'",
    ]
    for i, item in enumerate(info_items):
        if i == 0:
            p = tf4b.paragraphs[0]
        else:
            p = tf4b.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.name = "Calibri"
        run.font.size = Pt(16)
        p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
