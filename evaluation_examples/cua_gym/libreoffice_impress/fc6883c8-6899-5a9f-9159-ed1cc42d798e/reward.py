"""
Reward Script: Animated organizational chart on slide 2
Task ID: impress_gf4_033
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): 10 rectangle shapes on slide 2 in 3-level hierarchy (CEO, 3 VP, 6 reports)
  Component 2 (0.15): Connector line shapes linking hierarchy levels
  Component 3 (0.25): Animations exist and use 'Appear' entrance effect (presetID=1)
  Component 4 (0.15): CEO box has click-triggered animation; VP boxes appear together after previous
  Component 5 (0.15): Report boxes appear in pairs with delays
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_033'

NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'


def persist_app_state(domain):
    """Best-effort save of any open LibreOffice documents."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_animations(pptx_path, slide_idx):
    """Parse animation data from slide XML. Returns list of dicts with spid, presetID, presetClass, nodeType, delay."""
    animations = []
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        try:
            with zf.open(f'ppt/slides/slide{slide_idx + 1}.xml') as f:
                root = ET.fromstring(f.read())
        except KeyError:
            return animations

    for ctn in root.iter(f'{{{NS_P}}}cTn'):
        preset_id = ctn.get('presetID')
        node_type = ctn.get('nodeType')
        if preset_id:
            sp_tgt = ctn.find(f'.//{{{NS_P}}}spTgt')
            spid = sp_tgt.get('spid') if sp_tgt is not None else None
            st_cond = ctn.find(f'{{{NS_P}}}stCondLst/{{{NS_P}}}cond')
            delay = st_cond.get('delay') if st_cond is not None else '0'
            animations.append({
                'spid': spid,
                'presetID': preset_id,
                'presetClass': ctn.get('presetClass'),
                'nodeType': node_type,
                'delay': delay,
            })
    return animations


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 2")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[1]  # Slide 2 (0-indexed)

    # Categorize shapes on slide 2
    rect_shapes = []
    line_shapes = []
    for shape in slide.shapes:
        stype = shape.shape_type
        # AUTO_SHAPE (1) = rectangles/rounded rectangles; also check FREEFORM (5)
        if stype == MSO_SHAPE_TYPE.AUTO_SHAPE:
            rect_shapes.append(shape)
        elif stype == MSO_SHAPE_TYPE.LINE:
            line_shapes.append(shape)

    # ------------------------------------------------------------------
    # Component 1: 10 rectangle shapes in 3-level hierarchy (0.30 points)
    # ------------------------------------------------------------------
    try:
        num_rects = len(rect_shapes)
        print(f"  Found {num_rects} rectangle/auto shapes on slide 2")

        # Check for 3-level hierarchy by analyzing vertical positions
        # Collect top positions of rectangles
        tops = sorted(set(s.top for s in rect_shapes))

        # Group by approximate vertical level (tolerance 10% of slide height)
        # Slide height ~ 6858000 EMU (7.5 inches)
        level_tolerance = 500000  # ~0.55 inches tolerance
        levels = []
        for t in tops:
            placed = False
            for lvl in levels:
                if abs(t - lvl[0]) < level_tolerance:
                    lvl.append(t)
                    placed = True
                    break
            if not placed:
                levels.append([t])

        num_levels = len(levels)
        print(f"  Found {num_levels} vertical levels in rectangle layout")

        # Score: need 10 rects and 3 levels
        if num_rects >= 10 and num_levels >= 3:
            print(f"PASS: Component 1 — 10+ rectangles in 3+ levels ({num_rects} rects, {num_levels} levels) (0.30 pts)")
            total_score += 0.30
        elif num_rects >= 8 and num_levels >= 3:
            print(f"PARTIAL: Component 1 — {num_rects} rects (need 10), {num_levels} levels (0.15 pts)")
            total_score += 0.15
        elif num_rects >= 4 and num_levels >= 2:
            print(f"PARTIAL: Component 1 — {num_rects} rects, {num_levels} levels (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 1 — only {num_rects} rects and {num_levels} levels")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # ------------------------------------------------------------------
    # Component 2: Connector lines linking hierarchy levels (0.15 points)
    # ------------------------------------------------------------------
    try:
        num_lines = len(line_shapes)
        print(f"  Found {num_lines} line/connector shapes on slide 2")

        # Need at least 9 connectors: 3 (CEO->VPs) + 6 (VPs->reports)
        if num_lines >= 9:
            print(f"PASS: Component 2 — {num_lines} connector lines (need >= 9) (0.15 pts)")
            total_score += 0.15
        elif num_lines >= 6:
            print(f"PARTIAL: Component 2 — {num_lines} connector lines (need >= 9) (0.08 pts)")
            total_score += 0.08
        elif num_lines >= 3:
            print(f"PARTIAL: Component 2 — {num_lines} connector lines (need >= 9) (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 2 — only {num_lines} connector lines")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # ------------------------------------------------------------------
    # Component 3: Animations use 'Appear' entrance effect (0.25 points)
    # ------------------------------------------------------------------
    try:
        animations = get_animations(file_path, 1)  # slide index 1
        print(f"  Found {len(animations)} animations on slide 2")

        if len(animations) == 0:
            print(f"FAIL: Component 3 — no animations found on slide 2")
        else:
            # Check that all animations are entrance type with presetID=1 (Appear)
            entrance_appear = [a for a in animations if a['presetClass'] == 'entr' and a['presetID'] == '1']
            entrance_other = [a for a in animations if a['presetClass'] == 'entr' and a['presetID'] != '1']
            non_entrance = [a for a in animations if a['presetClass'] != 'entr']

            print(f"  Appear entrance: {len(entrance_appear)}, other entrance: {len(entrance_other)}, non-entrance: {len(non_entrance)}")

            # Need animations for at least the 10 rect shapes
            animated_rect_spids = set()
            for shape in rect_shapes:
                for a in animations:
                    if a['spid'] == str(shape.shape_id):
                        animated_rect_spids.add(shape.shape_id)
                        break

            num_animated_rects = len(animated_rect_spids)
            print(f"  Animated rectangle shapes: {num_animated_rects}/{len(rect_shapes)}")

            # Score based on having Appear animations for rectangles
            if num_animated_rects >= 10 and len(entrance_appear) >= 10:
                print(f"PASS: Component 3 — all 10+ rects animated with Appear (0.25 pts)")
                total_score += 0.25
            elif num_animated_rects >= 7 and len(entrance_appear) >= 7:
                print(f"PARTIAL: Component 3 — {num_animated_rects} rects animated (0.15 pts)")
                total_score += 0.15
            elif num_animated_rects >= 4:
                print(f"PARTIAL: Component 3 — {num_animated_rects} rects animated (0.08 pts)")
                total_score += 0.08
            else:
                print(f"FAIL: Component 3 — only {num_animated_rects} rects have animations")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # ------------------------------------------------------------------
    # Component 4: CEO click-triggered, VPs appear together after previous (0.15 points)
    # ------------------------------------------------------------------
    try:
        animations = get_animations(file_path, 1)

        # Build spid -> animation info map
        anim_by_spid = {}
        for a in animations:
            spid = a['spid']
            if spid not in anim_by_spid:
                anim_by_spid[spid] = a

        # Find CEO shape: should be at the top level (smallest top value among rects)
        if len(rect_shapes) >= 4:
            # Sort rects by top position
            sorted_rects = sorted(rect_shapes, key=lambda s: s.top)
            ceo_shape = sorted_rects[0]
            ceo_spid = str(ceo_shape.shape_id)

            # VP shapes: next 3 by top position (second level)
            # Group by approximate level
            level_groups = {}
            for s in rect_shapes:
                placed = False
                for lvl_top in level_groups:
                    if abs(s.top - lvl_top) < 500000:
                        level_groups[lvl_top].append(s)
                        placed = True
                        break
                if not placed:
                    level_groups[s.top] = [s]

            sorted_levels = sorted(level_groups.keys())
            score_4 = 0.0

            # CEO animation should be clickEffect
            if ceo_spid in anim_by_spid:
                ceo_anim = anim_by_spid[ceo_spid]
                if ceo_anim['nodeType'] == 'clickEffect':
                    print(f"  CEO (spid={ceo_spid}): clickEffect - correct")
                    score_4 += 0.075
                else:
                    print(f"  CEO (spid={ceo_spid}): {ceo_anim['nodeType']} - expected clickEffect")
            else:
                print(f"  CEO (spid={ceo_spid}): no animation found")

            # VP shapes should include afterEffect and withEffect (appear together)
            if len(sorted_levels) >= 2:
                vp_shapes = level_groups[sorted_levels[1]]
                vp_spids = [str(s.shape_id) for s in vp_shapes]
                vp_types = []
                for spid in vp_spids:
                    if spid in anim_by_spid:
                        vp_types.append(anim_by_spid[spid]['nodeType'])

                has_after = 'afterEffect' in vp_types
                has_with = 'withEffect' in vp_types
                print(f"  VP shapes (spids={vp_spids}): types={vp_types}")

                # At least one afterEffect and at least one withEffect means they appear together
                if has_after and has_with and len(vp_types) >= 3:
                    print(f"  VPs appear together (afterEffect + withEffect) - correct")
                    score_4 += 0.075
                elif has_after or has_with:
                    print(f"  VPs partially correct animation grouping")
                    score_4 += 0.04
                else:
                    print(f"  VPs missing correct animation grouping")

            if score_4 > 0:
                print(f"PASS: Component 4 — CEO click + VP grouping ({score_4:.3f} pts)")
                total_score += score_4
            else:
                print(f"FAIL: Component 4 — animation sequence not correct")
        else:
            print(f"FAIL: Component 4 — insufficient rectangle shapes for hierarchy")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # ------------------------------------------------------------------
    # Component 5: Report boxes appear in pairs with delays (0.15 points)
    # ------------------------------------------------------------------
    try:
        animations = get_animations(file_path, 1)

        # Build spid -> animation info map
        anim_by_spid = {}
        for a in animations:
            spid = a['spid']
            if spid not in anim_by_spid:
                anim_by_spid[spid] = a

        if len(rect_shapes) >= 4:
            # Identify bottom-level shapes (reports)
            level_groups = {}
            for s in rect_shapes:
                placed = False
                for lvl_top in level_groups:
                    if abs(s.top - lvl_top) < 500000:
                        level_groups[lvl_top].append(s)
                        placed = True
                        break
                if not placed:
                    level_groups[s.top] = [s]

            sorted_levels = sorted(level_groups.keys())

            if len(sorted_levels) >= 3:
                report_shapes = level_groups[sorted_levels[2]]
                report_spids = [str(s.shape_id) for s in report_shapes]
                print(f"  Report shapes (bottom level): spids={report_spids}")

                # Check for delayed animations (afterEffect with delay > 0)
                delayed_anims = []
                for spid in report_spids:
                    if spid in anim_by_spid:
                        a = anim_by_spid[spid]
                        if a['nodeType'] == 'afterEffect' and a['delay'] != '0':
                            delayed_anims.append(spid)

                # Check for pair grouping (withEffect among reports)
                with_effect_reports = []
                for spid in report_spids:
                    if spid in anim_by_spid:
                        a = anim_by_spid[spid]
                        if a['nodeType'] == 'withEffect':
                            with_effect_reports.append(spid)

                total_report_anims = sum(1 for spid in report_spids if spid in anim_by_spid)
                print(f"  Reports: {total_report_anims} animated, {len(delayed_anims)} with delay, {len(with_effect_reports)} withEffect")

                # Expect 3 pairs: each pair has 1 afterEffect with delay + 1 withEffect
                if len(delayed_anims) >= 3 and len(with_effect_reports) >= 3 and total_report_anims >= 6:
                    print(f"PASS: Component 5 — report pairs with delays (0.15 pts)")
                    total_score += 0.15
                elif len(delayed_anims) >= 2 and total_report_anims >= 4:
                    print(f"PARTIAL: Component 5 — some report pairs with delays (0.08 pts)")
                    total_score += 0.08
                elif total_report_anims >= 3:
                    print(f"PARTIAL: Component 5 — reports have animations but grouping incomplete (0.05 pts)")
                    total_score += 0.05
                else:
                    print(f"FAIL: Component 5 — report animations insufficient")
            else:
                print(f"FAIL: Component 5 — fewer than 3 levels found")
        else:
            print(f"FAIL: Component 5 — insufficient shapes")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist app state before verification
persist_app_state("libreoffice_impress")

# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
