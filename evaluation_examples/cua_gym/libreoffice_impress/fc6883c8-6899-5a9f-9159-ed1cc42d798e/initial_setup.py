"""
Initial Setup: 6-slide presentation with blank org chart slide
Task ID: impress_gf4_033
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_033'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_shape(shape, text, font_size=Pt(18), bold=False, color=None, alignment=None):
    """Helper to set text on a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title Slide ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Nexus Technologies Inc."
    slide1.placeholders[1].text = "Annual Corporate Overview 2025"

    # ── Slide 2: Organizational Chart (blank content area) ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title text box at the top
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Organizational Chart"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x6E)
    # No other shapes, no animations - blank content area

    # ── Slide 3: Company Mission ──
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    tb3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Our Mission & Vision"
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.runs[0]
    r3.font.size = Pt(32)
    r3.font.bold = True
    r3.font.color.rgb = RGBColor(0x2E, 0x4A, 0x6E)

    content3 = slide3.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5))
    tf3c = content3.text_frame
    tf3c.word_wrap = True
    p = tf3c.paragraphs[0]
    p.text = "At Nexus Technologies, we are committed to delivering innovative solutions that transform how enterprises manage their digital infrastructure. Our mission is to empower organizations with cutting-edge technology platforms."
    p.runs[0].font.size = Pt(18)

    p2 = tf3c.add_paragraph()
    p2.text = ""
    p3a = tf3c.add_paragraph()
    p3a.text = "Founded in 2012, Nexus Technologies has grown from a startup of 15 engineers to a global enterprise with over 2,400 employees across 8 countries. Our core products serve Fortune 500 companies in healthcare, finance, and manufacturing."
    p3a.runs[0].font.size = Pt(18)

    # ── Slide 4: Revenue Highlights ──
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    tb4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf4 = tb4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Revenue Highlights Q1-Q4 2024"
    p4.alignment = PP_ALIGN.CENTER
    r4 = p4.runs[0]
    r4.font.size = Pt(32)
    r4.font.bold = True
    r4.font.color.rgb = RGBColor(0x2E, 0x4A, 0x6E)

    # Add a table for revenue data
    table_shape = slide4.shapes.add_table(5, 4, Inches(1.5), Inches(1.5), Inches(10), Inches(3.5))
    table = table_shape.table
    headers = ["Quarter", "Revenue ($M)", "Growth (%)", "Net Margin (%)"]
    data = [
        ["Q1 2024", "$142.3", "12.4%", "18.7%"],
        ["Q2 2024", "$158.7", "14.2%", "19.3%"],
        ["Q3 2024", "$171.5", "16.8%", "20.1%"],
        ["Q4 2024", "$189.2", "18.1%", "21.5%"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # ── Slide 5: Strategic Priorities ──
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    tb5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf5 = tb5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Strategic Priorities 2025"
    p5.alignment = PP_ALIGN.CENTER
    r5 = p5.runs[0]
    r5.font.size = Pt(32)
    r5.font.bold = True
    r5.font.color.rgb = RGBColor(0x2E, 0x4A, 0x6E)

    priorities = [
        "Expand AI-powered analytics suite to 3 new verticals",
        "Achieve SOC 2 Type II and ISO 27001 certifications",
        "Launch European data center in Frankfurt by Q3",
        "Grow enterprise customer base by 25% year-over-year",
        "Increase R&D investment to 22% of revenue",
    ]
    content5 = slide5.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5))
    tf5c = content5.text_frame
    tf5c.word_wrap = True
    for i, pri in enumerate(priorities):
        if i == 0:
            p = tf5c.paragraphs[0]
        else:
            p = tf5c.add_paragraph()
        p.text = f"  {i+1}. {pri}"
        p.runs[0].font.size = Pt(18)
        p.space_after = Pt(12)

    # ── Slide 6: Contact Information ──
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    tb6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf6 = tb6.text_frame
    p6 = tf6.paragraphs[0]
    p6.text = "Contact Us"
    p6.alignment = PP_ALIGN.CENTER
    r6 = p6.runs[0]
    r6.font.size = Pt(32)
    r6.font.bold = True
    r6.font.color.rgb = RGBColor(0x2E, 0x4A, 0x6E)

    contact_info = slide6.shapes.add_textbox(Inches(3), Inches(2), Inches(7), Inches(4))
    tf6c = contact_info.text_frame
    tf6c.word_wrap = True
    lines = [
        "Nexus Technologies Inc.",
        "1200 Innovation Boulevard, Suite 400",
        "San Jose, CA 95134",
        "",
        "General Inquiries: info@nexustech.com",
        "Investor Relations: ir@nexustech.com",
        "Phone: +1 (408) 555-0192",
    ]
    for i, line in enumerate(lines):
        if i == 0:
            p = tf6c.paragraphs[0]
        else:
            p = tf6c.add_paragraph()
        p.text = line
        if line:
            p.runs[0].font.size = Pt(18)
            if i == 0:
                p.runs[0].font.bold = True
                p.runs[0].font.size = Pt(22)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
