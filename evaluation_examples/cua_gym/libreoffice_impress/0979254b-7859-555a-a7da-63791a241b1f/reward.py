"""
Reward Script: Health and Wellness Workshop Presentation
Task ID: impress_wf_042
Domain: libreoffice_impress
Scoring:
  C1: File exists + exactly 10 slides (0.15)
  C2: Slide 1 title contains "Wellness at Work" (0.10)
  C3: Slide 2 has hyperlinks linking to other slides (0.15)
  C4: Slide 4 has stick-figure shapes (oval + rectangles/lines) (0.10)
  C5: Slide 5 has pie/arc shapes for food plate diagram (0.15)
  C6: Slide 7 has a table (mental health resources) (0.10)
  C7: Slide 8 has a 6x5 table (30-day challenge calendar) (0.10)
  C8: Slide 9 has multiple text boxes (goal-setting worksheet) (0.05)
  C9: Green #81C784 is used as accent color (0.10)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_042'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Try to load file
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print("CRITICAL: Cannot load file %s: %s" % (file_path, e))
        print("REWARD: 0.0")
        return 0.0

    # Component 1: File has exactly 10 slides (0.15 points)
    try:
        num_slides = len(prs.slides)
        if num_slides == 10:
            print("PASS: Component 1 — 10 slides found (0.15 pts)")
            total_score += 0.15
        else:
            print("FAIL: Component 1 — expected 10 slides, found %d" % num_slides)
    except Exception as e:
        print("ERROR: Component 1 — %s" % e)

    # Component 2: Slide 1 title contains "Wellness at Work" (0.10 points)
    try:
        slide1 = prs.slides[0]
        slide1_text = ""
        for shape in slide1.shapes:
            if shape.has_text_frame:
                slide1_text += " " + shape.text_frame.text
        if "wellness at work" in slide1_text.lower():
            print("PASS: Component 2 — Slide 1 has 'Wellness at Work' title (0.10 pts)")
            total_score += 0.10
        else:
            print("FAIL: Component 2 — 'Wellness at Work' not found in slide 1 text: %r" % slide1_text[:200])
    except Exception as e:
        print("ERROR: Component 2 — %s" % e)

    # Component 3: Slide 2 has hyperlinks to other slides (0.15 points)
    # We check via XML for hlinkClick elements with action ppaction://hlinksldjump
    try:
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        hyperlink_count = 0
        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide2.xml') as f:
                root = ET.parse(f).getroot()
                hlinks = root.findall('.//{%s}hlinkClick' % ns_a)
                for h in hlinks:
                    action = h.get('action', '')
                    if 'hlinksldjump' in action:
                        hyperlink_count += 1
        if hyperlink_count >= 3:
            print("PASS: Component 3 — Slide 2 has %d slide hyperlinks (0.15 pts)" % hyperlink_count)
            total_score += 0.15
        else:
            print("FAIL: Component 3 — expected >=3 slide hyperlinks on slide 2, found %d" % hyperlink_count)
    except Exception as e:
        print("ERROR: Component 3 — %s" % e)

    # Component 4: Slide 4 has stick-figure shapes — oval (head) + rectangles/lines (body) (0.10 points)
    try:
        slide4 = prs.slides[3]
        has_oval = False
        rect_line_count = 0
        for shape in slide4.shapes:
            shape_name_lower = shape.name.lower()
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if 'oval' in shape_name_lower or 'circle' in shape_name_lower or 'ellipse' in shape_name_lower:
                    has_oval = True
                elif 'rectangle' in shape_name_lower or 'line' in shape_name_lower:
                    rect_line_count += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
                rect_line_count += 1
        if has_oval and rect_line_count >= 2:
            print("PASS: Component 4 — Slide 4 has oval + %d rect/line shapes for stick figure (0.10 pts)" % rect_line_count)
            total_score += 0.10
        else:
            print("FAIL: Component 4 — Slide 4 stick figure: oval=%s, rect/lines=%d (need oval + >=2)" % (has_oval, rect_line_count))
    except Exception as e:
        print("ERROR: Component 4 — %s" % e)

    # Component 5: Slide 5 has pie/arc shapes for food plate diagram (0.15 points)
    try:
        slide5 = prs.slides[4]
        pie_arc_count = 0
        for shape in slide5.shapes:
            shape_name_lower = shape.name.lower()
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if 'pie' in shape_name_lower or 'arc' in shape_name_lower or 'chord' in shape_name_lower:
                    pie_arc_count += 1
        if pie_arc_count >= 3:
            print("PASS: Component 5 — Slide 5 has %d pie/arc shapes for food plate (0.15 pts)" % pie_arc_count)
            total_score += 0.15
        else:
            print("FAIL: Component 5 — expected >=3 pie/arc shapes on slide 5, found %d" % pie_arc_count)
    except Exception as e:
        print("ERROR: Component 5 — %s" % e)

    # Component 6: Slide 7 has a table (mental health resources) (0.10 points)
    try:
        slide7 = prs.slides[6]
        table_found = False
        for shape in slide7.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                t = shape.table
                if len(t.rows) >= 2 and len(t.columns) >= 2:
                    table_found = True
                    print("PASS: Component 6 — Slide 7 has a %dx%d table (0.10 pts)" % (len(t.rows), len(t.columns)))
                    total_score += 0.10
                    break
        if not table_found:
            print("FAIL: Component 6 — No table with >=2 rows and >=2 cols found on slide 7")
    except Exception as e:
        print("ERROR: Component 6 — %s" % e)

    # Component 7: Slide 8 has a 6x5 table (30-day challenge calendar) (0.10 points)
    try:
        slide8 = prs.slides[7]
        calendar_found = False
        for shape in slide8.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                t = shape.table
                rows = len(t.rows)
                cols = len(t.columns)
                # Accept 6x5 or 5x6 (rows x cols either way) or close variations
                if (rows == 6 and cols == 5) or (rows == 5 and cols == 6):
                    calendar_found = True
                    print("PASS: Component 7 — Slide 8 has %dx%d calendar table (0.10 pts)" % (rows, cols))
                    total_score += 0.10
                    break
                # Also accept if total cells >= 30 (30-day calendar)
                elif rows * cols >= 30:
                    calendar_found = True
                    print("PASS: Component 7 — Slide 8 has %dx%d table with %d cells (0.10 pts)" % (rows, cols, rows * cols))
                    total_score += 0.10
                    break
        if not calendar_found:
            print("FAIL: Component 7 — No 6x5 table found on slide 8")
    except Exception as e:
        print("ERROR: Component 7 — %s" % e)

    # Component 8: Slide 9 has multiple text boxes for goal-setting worksheet (0.05 points)
    try:
        slide9 = prs.slides[8]
        textbox_count = 0
        for shape in slide9.shapes:
            if shape.has_text_frame:
                textbox_count += 1
        if textbox_count >= 5:
            print("PASS: Component 8 — Slide 9 has %d text boxes for worksheet (0.05 pts)" % textbox_count)
            total_score += 0.05
        else:
            print("FAIL: Component 8 — expected >=5 text boxes on slide 9, found %d" % textbox_count)
    except Exception as e:
        print("ERROR: Component 8 — %s" % e)

    # Component 9: Green #81C784 is used as accent color (0.10 points)
    # Check background fills and shape fills across all slides for the green color
    try:
        green_found = False
        green_hex = '81C784'
        for slide in prs.slides:
            # Check background
            try:
                bg = slide.background.fill
                if bg.type is not None and bg.type == 1:
                    rgb = str(bg.fore_color.rgb).upper()
                    if green_hex in rgb:
                        green_found = True
                        break
            except:
                pass
            # Check shape fills
            for shape in slide.shapes:
                try:
                    if hasattr(shape, 'fill') and shape.fill.type is not None and shape.fill.type == 1:
                        rgb = str(shape.fill.fore_color.rgb).upper()
                        if green_hex in rgb:
                            green_found = True
                            break
                except:
                    pass
                # Check text color
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                if run.font.color.type is not None:
                                    rgb = str(run.font.color.rgb).upper()
                                    if green_hex in rgb:
                                        green_found = True
                                        break
                            except:
                                pass
                        if green_found:
                            break
                if green_found:
                    break
            if green_found:
                break

        if green_found:
            print("PASS: Component 9 — Green #81C784 accent color found (0.10 pts)")
            total_score += 0.10
        else:
            print("FAIL: Component 9 — Green #81C784 not found in any background, shape fill, or text color")
    except Exception as e:
        print("ERROR: Component 9 — %s" % e)

    final_score = min(total_score, 1.0)
    print("")
    print("Score: %.2f/1.0" % total_score)
    print("REWARD: %.1f" % final_score)
    return final_score


# Persist app state before verification (in case GUI edits are unsaved)
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print("PERSIST_WARN: save hook failed: %s" % e)


# Entry point
file_path = os.path.join(WORKDIR, 'Desktop', 'Wellness_Workshop.pptx')
if not os.path.exists(file_path):
    print("File not found: %s" % file_path)
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
