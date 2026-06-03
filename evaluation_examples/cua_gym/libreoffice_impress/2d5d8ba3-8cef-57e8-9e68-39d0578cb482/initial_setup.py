"""
Initial Setup: Create presentation with column chart (no error bars)
Task ID: impress_tct_066
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_066'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Error Margins Analysis"
    slide1.placeholders[1].text = "Q4 2025 Regional Performance Review"

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Overview"
    tf = slide2.placeholders[1].text_frame
    tf.text = "This report examines regional sales performance for Q4 2025."
    p = tf.add_paragraph()
    p.text = "Key findings include variance analysis across six territories."
    p.level = 0
    p2 = tf.add_paragraph()
    p2.text = "Error margins highlight the uncertainty in projected versus actual results."
    p2.level = 0

    # --- Slide 3: Column Chart (NO error bars) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    # Add title text box
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf3 = txBox.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Regional Sales (Thousands USD)"
    p3.alignment = PP_ALIGN.CENTER
    run = p3.runs[0]
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2F, 0x54, 0x96)

    # Add column chart with 1 data series, 6 data points
    chart_data = CategoryChartData()
    chart_data.categories = [
        'Northeast', 'Southeast', 'Midwest',
        'Southwest', 'Northwest', 'Central'
    ]
    chart_data.add_series('Q4 Sales', (45, 32, 28, 51, 37, 42))

    chart_frame = slide3.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8), Inches(1.0),
        Inches(8.4), Inches(5.5),
        chart_data
    )

    chart = chart_frame.chart
    chart.has_legend = True
    chart.style = 2

    # Style the series
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = RGBColor(0x4E, 0x79, 0xA7)

    # --- Slide 4: Summary ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Summary & Next Steps"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "The error margin analysis reveals consistent performance patterns."
    p4a = tf4.add_paragraph()
    p4a.text = "Recommended actions:"
    p4a.level = 0
    p4b = tf4.add_paragraph()
    p4b.text = "Increase investment in Southwest and Northeast regions."
    p4b.level = 1
    p4c = tf4.add_paragraph()
    p4c.text = "Monitor Central region for emerging growth trends."
    p4c.level = 1
    p4d = tf4.add_paragraph()
    p4d.text = "Target Q1 2026 review for updated projections."
    p4d.level = 1

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
