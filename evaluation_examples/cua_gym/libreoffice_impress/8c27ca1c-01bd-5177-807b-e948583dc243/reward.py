"""
Reward Script: Update slide 5 background to amber (#FFBF00) and add speaker note
Task ID: osworld_impress_note_bg_combined_006
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): Slide 5 background is amber #FFBF00
  Component 2 (0.5): Slide 5 speaker notes contain the required text
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_note_bg_combined_006'

# Expected values from task context
EXPECTED_BG_COLOR = 'FFBF00'
EXPECTED_NOTE = 'This is the pivotal slide \u2014 make sure to pause and let the data sink in'
TARGET_SLIDE_IDX = 4  # 0-based index for slide 5


def get_slide_background_rgb(slide):
    """Return the hex string of the slide's solid background fill, or None."""
    fill = slide.background.fill
    if fill.type == 1:  # SOLID
        return str(fill.fore_color.rgb)
    elif fill.type == 5:  # Inherited from master
        master_fill = slide.slide_layout.slide_master.background.fill
        if master_fill.type == 1:
            return str(master_fill.fore_color.rgb)
    return None


def get_slide_notes(slide):
    """Return the speaker notes text of a slide, stripped."""
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ''


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition check: file must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"CRITICAL: Expected at least 5 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide5 = prs.slides[TARGET_SLIDE_IDX]

    # Component 1: Slide 5 background is amber #FFBF00 (0.5 points)
    # Initial env has FFFFFF; golden env should have FFBF00.
    try:
        actual_bg = get_slide_background_rgb(slide5)
        if actual_bg is not None and actual_bg.upper() == EXPECTED_BG_COLOR.upper():
            print(f"PASS: Component 1 — Slide 5 background is #{actual_bg} (expected #{EXPECTED_BG_COLOR}) (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 1 — Slide 5 background is #{actual_bg}, expected #{EXPECTED_BG_COLOR}")
    except Exception as e:
        print(f"ERROR: Component 1 — Could not check slide 5 background: {e}")

    # Component 2: Slide 5 speaker notes contain the required text (0.5 points)
    # Initial env has empty notes; golden env should have the exact required note.
    try:
        actual_notes = get_slide_notes(slide5)
        if actual_notes == EXPECTED_NOTE:
            print(f"PASS: Component 2 — Slide 5 notes match expected text (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 2 — Slide 5 notes: {repr(actual_notes)}")
            print(f"      Expected:                    {repr(EXPECTED_NOTE)}")
    except Exception as e:
        print(f"ERROR: Component 2 — Could not check slide 5 notes: {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the VM environment
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
