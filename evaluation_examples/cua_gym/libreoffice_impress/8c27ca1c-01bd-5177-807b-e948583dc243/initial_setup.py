"""
Initial Setup: Investor pitch deck — 6 slides, slide 5 has white background and no notes
Task ID: osworld_impress_note_bg_combined_006
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_note_bg_combined_006'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_slide_background_white(slide):
    """Explicitly set slide background to white."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_text(slide, title_text, subtitle_text=None):
    """Add title and optional subtitle to slide using text boxes."""
    # Title text box
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.name = "Calibri"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)

    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(0.6))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle_text
        run2.font.name = "Calibri"
        run2.font.size = Pt(16)
        run2.font.color.rgb = RGBColor(0x44, 0x47, 0x2A)


def add_body_text(slide, lines, top_inches=2.4):
    """Add bullet-point style body text to a slide."""
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(top_inches), Inches(8.6), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = "Calibri"
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def create_initial():
    prs = Presentation()
    # Use 16:9 widescreen dimensions
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    # ---- Slide 1: Company Overview ----
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide1)
    add_title_text(slide1, "NovaTech Solutions", "Series B Investor Presentation — Q1 2026")
    add_body_text(slide1, [
        "Founded: 2019  |  Headquarters: San Francisco, CA",
        "Industry: Enterprise SaaS / AI-Powered Analytics",
        "Team: 142 employees across 4 global offices",
        "Mission: Empower data-driven decisions for mid-market enterprises",
    ])

    # ---- Slide 2: Problem Statement ----
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide2)
    add_title_text(slide2, "The Problem", "Mid-market businesses are drowning in data noise")
    add_body_text(slide2, [
        "73% of enterprise data goes unanalyzed (Gartner 2025)",
        "Average analyst spends 62% of time on data prep, not insights",
        "Legacy BI tools require months of implementation and specialist staff",
        "Result: Missed opportunities worth $2.3T annually across US mid-market",
    ])

    # ---- Slide 3: Our Solution ----
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide3)
    add_title_text(slide3, "Our Solution", "NovaTech Insight Engine™")
    add_body_text(slide3, [
        "One-click integration with 200+ data sources (CRM, ERP, cloud storage)",
        "AI-driven anomaly detection and predictive forecasting",
        "No-code dashboard builder — live in 48 hours, not 6 months",
        "Natural language query interface powered by proprietary LLM fine-tuning",
        "SOC 2 Type II certified; GDPR and CCPA compliant",
    ])

    # ---- Slide 4: Market Opportunity ----
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide4)
    add_title_text(slide4, "Market Opportunity", "A $47B and growing addressable market")
    add_body_text(slide4, [
        "Total Addressable Market (TAM): $47.2B by 2028",
        "Serviceable Addressable Market (SAM): $12.8B — US & Canada mid-market",
        "Serviceable Obtainable Market (SOM): $640M — 3-year target",
        "CAGR: 18.4% (Business Intelligence & Analytics sector)",
        "Key verticals: Healthcare, Retail, Manufacturing, Financial Services",
    ])

    # ---- Slide 5: Financial Projections (WHITE BG, NO NOTES — pre-task state) ----
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide5)
    add_title_text(slide5, "Financial Projections", "Path to profitability by Q3 2027")
    add_body_text(slide5, [
        "FY2024 ARR: $8.4M  |  Growth: +127% YoY",
        "FY2025 ARR Target: $19.2M  |  Projected Growth: +129% YoY",
        "FY2026 ARR Target: $41.7M  |  Projected Growth: +117% YoY",
        "Gross Margin: 74%  |  Net Revenue Retention: 118%",
        "Burn Rate: $1.1M/month  |  Runway: 18 months (pre-raise)",
        "Break-even: Q3 2027 at current growth trajectory",
    ])
    # NOTE: No speaker notes on slide 5 — this is the pre-task state

    # ---- Slide 6: Call to Action / Ask ----
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide6)
    add_title_text(slide6, "The Ask", "Join us in building the future of enterprise analytics")
    add_body_text(slide6, [
        "Raising: $18M Series B at $95M pre-money valuation",
        "Use of funds: 60% product R&D, 25% sales & marketing, 15% G&A",
        "Lead investors: Benchmark Capital (committed), Sequoia (in diligence)",
        "Close date target: April 30, 2026",
        "Contact: investors@novatech.io  |  www.novatech.io/investors",
    ])
    # Add notes to a non-task slide to ensure realistic content elsewhere
    slide6.notes_slide.notes_text_frame.text = "Thank the investors for their time. Remind them of the data room access link."

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
