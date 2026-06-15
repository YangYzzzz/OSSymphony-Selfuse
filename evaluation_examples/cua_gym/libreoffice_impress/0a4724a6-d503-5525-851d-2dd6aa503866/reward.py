"""
Reward Script: Agile Sprint Cycle Process Diagram on Slide 4
Task ID: impress_stu_089
Domain: libreoffice_impress
Scoring:
  C1 (0.30) - 5 main step rounded rectangles with correct labels
  C2 (0.15) - Main step shapes have light blue fill #D4E6F1
  C3 (0.15) - Daily Standup oval with light green fill #D5F5E3
  C4 (0.15) - Arrow connectors present (>= 5)
  C5 (0.25) - Center label 'Sprint (2 weeks)' in 20pt bold
"""

import os
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_089'

# The expected main step labels (normalized lowercase for flexible matching)
EXPECTED_MAIN_STEPS = {
    'backlog grooming',
    'sprint planning',
    'development',
    'sprint review',
    'retrospective',
}


def get_shape_geometry(shape):
    """Get the preset geometry type from shape XML."""
    try:
        sp_elem = shape._element
        ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        prstGeom = sp_elem.find(f'.//{ns}prstGeom')
        if prstGeom is not None:
            return prstGeom.get('prst')
    except Exception:
        pass
    return None


def get_shape_fill_color(shape):
    """Get the solid fill srgbClr hex value from shape XML."""
    try:
        ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        sp_elem = shape._element
        # Look for solidFill under spPr (shape properties)
        spPr = sp_elem.find(f'.//{ns}solidFill')
        if spPr is not None:
            srgb = spPr.find(f'{ns}srgbClr')
            if srgb is not None:
                return srgb.get('val').upper()
    except Exception:
        pass
    return None


def get_shape_text_normalized(shape):
    """Get shape text with vertical tabs replaced by spaces, stripped and lowered."""
    if hasattr(shape, 'text') and shape.text:
        return shape.text.replace('\x0b', ' ').strip().lower()
    return ''


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)

    # Categorize shapes on slide 4
    rounded_rects = []  # (text_lower, fill_color)
    ovals = []  # (text_lower, fill_color)
    connectors = []
    center_label_shapes = []

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            geom = get_shape_geometry(shape)
            text = get_shape_text_normalized(shape)
            fill = get_shape_fill_color(shape)
            if geom == 'roundRect':
                rounded_rects.append((text, fill))
            elif geom == 'ellipse':
                ovals.append((text, fill))
        elif shape.shape_type == 9:  # LINE / connector
            connectors.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            text_raw = shape.text.strip() if hasattr(shape, 'text') else ''
            if 'sprint' in text_raw.lower() and 'week' in text_raw.lower():
                center_label_shapes.append(shape)

    # ---- Component 1: 5 main step rounded rectangles with correct labels (0.30 pts) ----
    try:
        found_labels = {text for text, _ in rounded_rects}
        matched = EXPECTED_MAIN_STEPS.intersection(found_labels)
        match_count = len(matched)
        if match_count == 5:
            print(f"PASS: Component 1 -- All 5 main step rounded rectangles found: {matched} (0.30 pts)")
            total_score += 0.30
        elif match_count >= 3:
            partial = round(0.30 * match_count / 5, 2)
            print(f"PARTIAL: Component 1 -- {match_count}/5 main steps found: {matched} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- Only {match_count}/5 main steps found as rounded rectangles. "
                  f"Found texts: {[t for t, _ in rounded_rects]}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # ---- Component 2: Main step shapes have light blue fill #D4E6F1 (0.15 pts) ----
    try:
        if len(rounded_rects) == 0:
            print("FAIL: Component 2 -- No rounded rectangles found to check fill color")
        else:
            correct_fill_count = sum(1 for _, fill in rounded_rects if fill == 'D4E6F1')
            total_rects = len(rounded_rects)
            if correct_fill_count >= 5:
                print(f"PASS: Component 2 -- All main step shapes have fill #D4E6F1 (0.15 pts)")
                total_score += 0.15
            elif correct_fill_count >= 3:
                partial = round(0.15 * correct_fill_count / 5, 2)
                print(f"PARTIAL: Component 2 -- {correct_fill_count}/{total_rects} shapes have correct fill ({partial} pts)")
                total_score += partial
            else:
                fills = [fill for _, fill in rounded_rects]
                print(f"FAIL: Component 2 -- Only {correct_fill_count}/{total_rects} shapes have fill #D4E6F1. Found fills: {fills}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # ---- Component 3: Daily Standup oval with light green fill #D5F5E3 (0.15 pts) ----
    try:
        standup_ovals = [(text, fill) for text, fill in ovals if 'standup' in text or 'stand-up' in text or 'stand up' in text]
        if len(standup_ovals) == 0:
            # Also check if 'daily' is in an oval
            standup_ovals = [(text, fill) for text, fill in ovals if 'daily' in text]

        if len(standup_ovals) > 0:
            text, fill = standup_ovals[0]
            if fill == 'D5F5E3':
                print(f"PASS: Component 3 -- Daily Standup oval found with fill #D5F5E3 (0.15 pts)")
                total_score += 0.15
            elif fill is not None:
                # Partial: oval exists but wrong color
                print(f"PARTIAL: Component 3 -- Daily Standup oval found but fill is #{fill}, expected #D5F5E3 (0.07 pts)")
                total_score += 0.07
            elif fill is None:
                # Partial: oval exists but no fill color detected
                print(f"PARTIAL: Component 3 -- Daily Standup oval found but no fill color set (0.07 pts)")
                total_score += 0.07
        else:
            print(f"FAIL: Component 3 -- No oval shape with 'Daily Standup' text found. Ovals: {ovals}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # ---- Component 4: Arrow connectors present (>= 5) (0.15 pts) ----
    try:
        num_connectors = len(connectors)
        if num_connectors >= 5:
            print(f"PASS: Component 4 -- {num_connectors} connectors found (>= 5 required) (0.15 pts)")
            total_score += 0.15
        elif num_connectors >= 3:
            partial = round(0.15 * num_connectors / 5, 2)
            print(f"PARTIAL: Component 4 -- {num_connectors} connectors found (need >= 5) ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 -- Only {num_connectors} connectors found, need >= 5")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # ---- Component 5: Center label 'Sprint (2 weeks)' in 20pt bold (0.25 pts) ----
    try:
        if len(center_label_shapes) == 0:
            print("FAIL: Component 5 -- No center label with 'Sprint' and 'weeks' found")
        else:
            shape = center_label_shapes[0]
            label_text = shape.text.strip()
            has_correct_text = 'sprint' in label_text.lower() and '2 week' in label_text.lower()

            # Check font properties via actual API calls
            bold_runs = [r for p in shape.text_frame.paragraphs for r in p.runs if r.font.bold is True]
            sized_runs = [r for p in shape.text_frame.paragraphs for r in p.runs
                          if r.font.size is not None and r.font.size == Pt(20)]
            is_bold = len(bold_runs) > 0
            is_20pt = len(sized_runs) > 0

            if has_correct_text and is_bold and is_20pt:
                print(f"PASS: Component 5 -- Center label '{label_text}' is 20pt bold (0.25 pts)")
                total_score += 0.25
            elif has_correct_text and (is_bold or is_20pt):
                print(f"PARTIAL: Component 5 -- Center label found, bold={is_bold}, 20pt={is_20pt} (0.12 pts)")
                total_score += 0.12
            elif has_correct_text:
                print(f"PARTIAL: Component 5 -- Center label text correct but not bold/20pt (0.08 pts)")
                total_score += 0.08
            else:
                print(f"FAIL: Component 5 -- Label text mismatch: '{label_text}'")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
