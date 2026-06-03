"""
Initial Setup: Create Agile_Methods presentation with 8 slides, slide 4 empty for diagram task
Task ID: impress_stu_089
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_089'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.paragraphs[0].text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
    return slide


def add_title_only_slide(prs, title_text):
    """Add a slide with only a title (body area empty)."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title as a text box at the top
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    add_title_slide(prs, "Agile Methods in Software Engineering",
                    "CS 4310 - Fall 2025\nProfessor Elena Vasquez")

    # --- Slide 2: What is Agile? ---
    add_content_slide(prs, "What is Agile?", [
        "Iterative and incremental approach to software development",
        "Emphasizes flexibility, collaboration, and customer feedback",
        "Responds to change over following a rigid plan",
        "Delivers working software in short cycles (sprints)",
        "Cross-functional teams own the full delivery process",
        "Originated from the Agile Manifesto (2001)"
    ])

    # --- Slide 3: Agile Manifesto ---
    add_content_slide(prs, "The Agile Manifesto", [
        "Individuals and interactions over processes and tools",
        "Working software over comprehensive documentation",
        "Customer collaboration over contract negotiation",
        "Responding to change over following a plan",
        "",
        "While there is value in the items on the right,",
        "we value the items on the left more."
    ])

    # --- Slide 4: Sprint Lifecycle (EMPTY - task target) ---
    add_title_only_slide(prs, "Sprint Lifecycle")

    # --- Slide 5: Scrum Roles ---
    add_content_slide(prs, "Scrum Roles", [
        "Product Owner: Defines and prioritizes the product backlog",
        "Scrum Master: Facilitates the process, removes impediments",
        "Development Team: Self-organizing, cross-functional (5-9 members)",
        "Stakeholders: Provide feedback during Sprint Reviews",
        "Each role has distinct responsibilities that prevent bottlenecks"
    ])

    # --- Slide 6: User Stories & Estimation ---
    add_content_slide(prs, "User Stories & Estimation", [
        "Format: 'As a [user], I want [feature], so that [benefit]'",
        "Story Points: Relative sizing using Fibonacci (1, 2, 3, 5, 8, 13, 21)",
        "Planning Poker: Team consensus on effort estimation",
        "Velocity: Average story points completed per sprint",
        "Burndown Charts: Track remaining work over time"
    ])

    # --- Slide 7: Agile Tools ---
    add_content_slide(prs, "Agile Tools & Frameworks", [
        "JIRA: Industry-standard project tracking and sprint boards",
        "Trello: Lightweight Kanban-style task management",
        "Azure DevOps: Microsoft's integrated ALM platform",
        "GitHub Projects: Built-in issue tracking with automation",
        "Confluence: Team documentation and knowledge sharing",
        "Miro: Virtual whiteboarding for remote ceremonies"
    ])

    # --- Slide 8: Summary & Next Steps ---
    add_content_slide(prs, "Summary & Next Steps", [
        "Agile enables adaptive, customer-focused development",
        "Scrum provides a structured framework within Agile principles",
        "Next lecture: Kanban vs. Scrum - choosing the right framework",
        "Assignment: Form teams of 4 and create your first product backlog",
        "Reading: Chapters 3-4 of 'Agile Estimating and Planning' by Mike Cohn"
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
