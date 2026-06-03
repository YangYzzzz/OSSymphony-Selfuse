"""
Reward Script: Bold the title text on slide 2
Task ID: osworld_impress_title_selective_formatting_001
Domain: libreoffice_impress
Scoring:
  Component 1 (0.7): Slide 2 title placeholder has all runs with bold=True
  Component 2 (0.3): Slide 2 title text is exactly 'Market Analysis' AND bold=True
                     (compound: content preserved + formatting applied correctly)
Both components require bold=True, so both fail on initial_env and pass on golden_env.
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_selective_formatting_001'

# Expected title text for slide 2 (from task context)
EXPECTED_SLIDE2_TITLE = 'Market Analysis'


def get_title_shape(slide):
    """Return the title placeholder shape from a slide, or None."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name.startswith('Title'):
            return shape
    return None


def is_run_bold(run):
    """
    Returns True only if run.font.bold is explicitly set to True.
    False/None both mean 'not bold'.
    """
    return run.font.bold is True


def nonempty_runs(para):
    """Return only runs with non-empty text."""
    return [r for r in para.runs if (r.text or "").strip()]


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition gate: must have exactly 5 slides
    if len(prs.slides) != 5:
        print(f"CRITICAL: Expected 5 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]  # 0-indexed: index 1 = slide 2

    # Locate the title shape on slide 2
    title_shape = get_title_shape(slide2)
    if title_shape is None:
        print("CRITICAL: No title shape found on slide 2")
        print("REWARD: 0.0")
        return 0.0

    # Gather all non-empty runs in the slide 2 title
    all_runs = []
    for para in title_shape.text_frame.paragraphs:
        all_runs.extend(nonempty_runs(para))

    title_text = title_shape.text_frame.text.strip()

    # Component 1: All runs in slide 2 title have bold=True (0.7 points)
    # FAILS on initial_env (bold=False) → PASSES on golden_env (bold=True)
    try:
        if not all_runs:
            print(f"FAIL: Component 1 — No non-empty runs in slide 2 title (text: {title_text!r})")
        else:
            all_bold = all(is_run_bold(r) for r in all_runs)
            if all_bold:
                bold_values = [(r.text, r.font.bold) for r in all_runs]
                print(f"PASS: Component 1 — All runs in slide 2 title are bold: {bold_values} (0.7 pts)")
                total_score += 0.7
            else:
                not_bold = [(r.text, r.font.bold) for r in all_runs if not is_run_bold(r)]
                print(f"FAIL: Component 1 — Some runs in slide 2 title are not bold: {not_bold}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 2 title text is 'Market Analysis' AND bold=True (0.3 points)
    # Compound check: verifies that (a) the bold formatting was applied AND (b) text content is intact.
    # FAILS on initial_env (bold=False, text matches but bold condition fails)
    # PASSES on golden_env (bold=True, text matches)
    try:
        title_text_matches = (title_text == EXPECTED_SLIDE2_TITLE)
        any_bold = any(is_run_bold(r) for r in all_runs) if all_runs else False

        if title_text_matches and any_bold:
            print(f"PASS: Component 2 — Title text is '{title_text}' and is bold (0.3 pts)")
            total_score += 0.3
        else:
            reasons = []
            if not title_text_matches:
                reasons.append(f"title text is '{title_text}', expected '{EXPECTED_SLIDE2_TITLE}'")
            if not any_bold:
                reasons.append("title runs are not bold")
            print(f"FAIL: Component 2 — {'; '.join(reasons)}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
