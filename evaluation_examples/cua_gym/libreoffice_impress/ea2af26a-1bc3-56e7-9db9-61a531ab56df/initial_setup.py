"""
Initial Setup: Bold the title text on slide 2
Task ID: osworld_impress_title_selective_formatting_001
Domain: libreoffice_impress

Creates a 5-slide business strategy deck.
Slide 2 title 'Market Analysis' is in Calibri, NOT bold.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_selective_formatting_001'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_title_text(slide, text, font_name='Calibri', font_size=36, bold=False,
                   color=None):
    """Helper to set title placeholder text with explicit font properties."""
    title_shape = slide.shapes.title
    tf = title_shape.text_frame
    tf.clear()
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)


def set_content_text(slide, text, placeholder_idx=1, font_name='Calibri',
                     font_size=18):
    """Helper to set content placeholder text."""
    try:
        ph = slide.placeholders[placeholder_idx]
        tf = ph.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(font_size)
    except (KeyError, IndexError):
        pass


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title slide ---
    layout_title = prs.slide_layouts[0]  # Title Slide
    slide1 = prs.slides.add_slide(layout_title)
    set_title_text(slide1, 'Q1 2025 Business Strategy', font_name='Calibri',
                   font_size=40, bold=False)
    try:
        slide1.placeholders[1].text = 'Strategic Planning Division\nMarch 2025'
    except (KeyError, IndexError):
        pass

    # --- Slide 2: Market Analysis (title MUST be non-bold) ---
    layout_content = prs.slide_layouts[1]  # Title + Content
    slide2 = prs.slides.add_slide(layout_content)
    set_title_text(slide2, 'Market Analysis', font_name='Calibri',
                   font_size=36, bold=False)  # NOT bold — task requires adding bold
    set_content_text(slide2,
                     'Total Addressable Market: $4.2B\n'
                     'Year-over-Year Growth: 18.3%\n'
                     'Key Segments: Enterprise (62%), SMB (28%), Consumer (10%)\n'
                     'Competitive Landscape: 3 major players hold 74% share\n'
                     'Emerging Trend: AI-driven analytics adoption up 45% QoQ',
                     placeholder_idx=1)

    # --- Slide 3: Revenue Performance ---
    slide3 = prs.slides.add_slide(layout_content)
    set_title_text(slide3, 'Revenue Performance', font_name='Calibri',
                   font_size=36, bold=False)
    set_content_text(slide3,
                     'Q1 Revenue: $12.7M (+22% YoY)\n'
                     'Recurring Revenue: $9.4M (74% of total)\n'
                     'New Bookings: $3.8M\n'
                     'Churn Rate: 4.1% (target: <5%)\n'
                     'Pipeline Coverage: 3.2x quota',
                     placeholder_idx=1)

    # --- Slide 4: Product Roadmap ---
    slide4 = prs.slides.add_slide(layout_content)
    set_title_text(slide4, 'Product Roadmap', font_name='Calibri',
                   font_size=36, bold=False)
    set_content_text(slide4,
                     'Q2 2025: Launch Analytics Dashboard v3.0\n'
                     'Q2 2025: Mobile app redesign (iOS & Android)\n'
                     'Q3 2025: API marketplace beta release\n'
                     'Q3 2025: Enterprise SSO integration\n'
                     'Q4 2025: AI-powered recommendation engine',
                     placeholder_idx=1)

    # --- Slide 5: Strategic Priorities ---
    slide5 = prs.slides.add_slide(layout_content)
    set_title_text(slide5, 'Strategic Priorities', font_name='Calibri',
                   font_size=36, bold=False)
    set_content_text(slide5,
                     'Priority 1: Expand enterprise customer base by 30%\n'
                     'Priority 2: Reduce customer acquisition cost by 15%\n'
                     'Priority 3: Achieve ISO 27001 certification by Q3\n'
                     'Priority 4: Launch partner ecosystem program\n'
                     'Priority 5: Improve NPS score from 42 to 55',
                     placeholder_idx=1)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
