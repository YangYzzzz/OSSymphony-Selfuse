"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 159 I’ve placed a rectangle, a circle, and a star one under the other, but the gaps look uneven. In LibreOffice Impress, how do I distribute those three shapes so the vertical spacing is exactly equal between each one?
Generated: 2025-09-10 17:36:28
Status: success
Model: azure-o3
Total Steps: 7
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE


def verify_equal_vertical_distribution(
    file_path: str,
    slide_idx_one_based: int = 159,
    tolerance_emu: int = 20_000,
) -> float:
    """
    Verify that on the specified slide the rectangle, circle and star shapes
    are present and that the vertical gaps between them are equal within a
    tolerance.  Uses progressive scoring:
        • 0.4 points – all three required shapes are present.
        • 0.6 points – the two gaps between the three shapes are equal within
          the specified tolerance.
    Returns a score between 0.0 and 1.0 (float).
    """

    print(f"Verifying vertical distribution on slide {slide_idx_one_based} in: {file_path}\n")

    score = 0.0
    max_score = 1.0

    # ------------------------------------------------------------------
    # 1) Load presentation ------------------------------------------------
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File does not exist")
        return 0.0  # No points for a missing file

    try:
        prs = Presentation(file_path)
    except Exception as exc:
        print(f"✗ Could not open presentation: {exc}")
        return 0.0  # Loading failure => task failed

    # ------------------------------------------------------------------
    # 2) Retrieve the targeted slide ------------------------------------
    # ------------------------------------------------------------------
    total_slides = len(prs.slides)
    if not (1 <= slide_idx_one_based <= total_slides):
        print(f"✗ Slide {slide_idx_one_based} does not exist (presentation has {total_slides} slides)")
        return 0.0

    slide = prs.slides[slide_idx_one_based - 1]

    # ------------------------------------------------------------------
    # 3) Locate the rectangle, circle, and star shapes -------------------
    # ------------------------------------------------------------------
    required_types = [
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        MSO_AUTO_SHAPE_TYPE.OVAL,        # circle/ellipse
        MSO_AUTO_SHAPE_TYPE.STAR_5_POINT,
    ]

    target_shapes = [
        shp
        for shp in slide.shapes
        if shp.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and shp.auto_shape_type in required_types
    ]

    print(f"Found {len(target_shapes)} target shapes (rectangle, circle, star)")

    # 3a) Check presence of each required shape (0.4 pts) --------------
    found_types = {shp.auto_shape_type for shp in target_shapes}

    if all(t in found_types for t in required_types):
        print("✓ All three required shapes are present (rectangle, circle, star)")
        score += 0.4
    else:
        missing = [t for t in required_types if t not in found_types]
        print("✗ Missing required shapes:", missing)

    # ------------------------------------------------------------------
    # 4) Verify equal vertical spacing (0.6 pts) -------------------------
    # ------------------------------------------------------------------
    if len(target_shapes) >= 3:
        # Sort shapes by their top coordinate (vertical position)
        ordered = sorted(target_shapes, key=lambda s: s.top)

        gaps = []  # gap[i] = space between shape[i] bottom & shape[i+1] top
        for idx in range(len(ordered) - 1):
            current_bottom = ordered[idx].top + ordered[idx].height
            next_top = ordered[idx + 1].top
            gap = next_top - current_bottom
            gaps.append(gap)
            print(f"Gap {idx + 1}: {gap} EMU")

        if len(gaps) >= 2:
            diff = abs(gaps[1] - gaps[0])
            if diff <= tolerance_emu:
                print(
                    f"✓ Gaps are equal within tolerance ({tolerance_emu} EMU). Difference: {diff}"
                )
                score += 0.6
            else:
                print(
                    f"✗ Gaps are not equal. Difference {diff} EMU exceeds tolerance {tolerance_emu}"
                )
        else:
            print("✗ Not enough gaps to verify spacing (need at least two)")
    else:
        print("✗ Less than three target shapes found; cannot verify spacing")

    # ------------------------------------------------------------------
    # 5) Finalise --------------------------------------------------------
    # ------------------------------------------------------------------
    final_score = min(score, max_score)
    print(f"Final score: {final_score}\n")
    return final_score


if __name__ == "__main__":
    FILE_PATH = (
        "/home/user/on_slide_159_ive_placed_a_rectangle_a_circle_and_a_star_one_"
        "under_the_other_but_the_gaps_look_uneven_golden.pptx"
    )

    reward = verify_equal_vertical_distribution(FILE_PATH)
    print(f"REWARD: {reward}")

