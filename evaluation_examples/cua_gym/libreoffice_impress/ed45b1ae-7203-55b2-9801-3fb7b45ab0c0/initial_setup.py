"""
Initial Setup: Create a 16:9 presentation with no guide lines
Task ID: impress_el_077
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Cm, Pt, Inches, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_el_077'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Set 16:9 dimensions: 33.867cm x 19.05cm
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q1 2025 Marketing Strategy"
    slide1.placeholders[1].text = "Prepared by the Brand & Digital Marketing Team"

    # --- Slide 2: Content slide with bullet points ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Campaign Objectives"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Increase brand awareness by 25% in the European market"
    p2 = body2.add_paragraph()
    p2.text = "Launch three targeted social media campaigns per month"
    p2.level = 0
    p3 = body2.add_paragraph()
    p3.text = "Achieve 15% growth in organic website traffic"
    p3.level = 0
    p4 = body2.add_paragraph()
    p4.text = "Expand influencer partnerships to 40 active collaborations"
    p4.level = 0
    p5 = body2.add_paragraph()
    p5.text = "Reduce customer acquisition cost by 12% through optimization"
    p5.level = 0

    # --- Slide 3: Another content slide ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Budget Allocation"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Digital Advertising: $185,000 (37%)"
    for item in [
        "Content Production: $95,000 (19%)",
        "Events & Sponsorships: $72,000 (14.4%)",
        "SEO & Analytics Tools: $48,000 (9.6%)",
        "Influencer Partnerships: $65,000 (13%)",
        "Contingency Reserve: $35,000 (7%)",
    ]:
        p = body3.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 4: Blank slide for visual content area ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide4.shapes.add_textbox(Cm(2), Cm(1), Cm(29), Cm(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Timeline & Key Milestones"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4E, 0x7E)

    # Add a text box with milestone details
    txBox2 = slide4.shapes.add_textbox(Cm(3), Cm(5), Cm(27), Cm(12))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    milestones = [
        "January 15 — Campaign brief finalized and approved",
        "February 1 — Creative assets production begins",
        "March 10 — Phase 1 digital ads go live across platforms",
        "April 5 — Mid-quarter performance review meeting",
        "May 20 — Phase 2 launch with adjusted targeting",
        "June 30 — End-of-quarter comprehensive report due",
    ]
    tf2.paragraphs[0].text = milestones[0]
    for m in milestones[1:]:
        p = tf2.add_paragraph()
        p.text = m

    # Save the presentation
    prs.save(OUTPUT)

    # Now remove all guide lines from viewProps.xml to ensure clean state
    # python-pptx creates default center guides; we need to remove them
    _remove_all_guides(OUTPUT)

    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


def _remove_all_guides(pptx_path):
    """Remove all guide lines from the presentation's viewProps.xml."""
    import shutil
    tmp_path = pptx_path + '.tmp'

    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ET.register_namespace('', ns_p)
    ET.register_namespace('a', 'http://schemas.openxmlformats.org/drawingml/2006/main')
    ET.register_namespace('r', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships')
    ET.register_namespace('p', ns_p)

    with zipfile.ZipFile(pptx_path, 'r') as zin, \
         zipfile.ZipFile(tmp_path, 'w') as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == 'ppt/viewProps.xml':
                root = ET.fromstring(data)
                # Find and remove all guideLst elements
                for parent in root.iter():
                    for child in list(parent):
                        if child.tag.endswith('}guideLst') or child.tag == 'guideLst':
                            parent.remove(child)
                data = b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n' + ET.tostring(root, encoding='unicode').encode('utf-8')
            zout.writestr(item, data)

    shutil.move(tmp_path, pptx_path)


create_initial()
