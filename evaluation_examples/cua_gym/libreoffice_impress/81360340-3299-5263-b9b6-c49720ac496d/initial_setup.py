"""
Initial Setup: Add exit animations to all objects on slide 9
Task ID: impress_fix_063
Domain: libreoffice_impress

Creates a 10-slide presentation 'Sequential_Display.pptx' where slide 9 has
4 objects (image, text box, chart-like shape, geometric shape) each with
'Appear' entrance animations triggered 'After Previous'. No exit animations.
"""

import os
import shlex
import subprocess
import time
import struct
import zlib

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import lxml.etree as etree

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_063'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_minimal_png(filepath, width=200, height=150, color=(70, 130, 180)):
    """Create a minimal valid PNG file without needing PIL."""
    def make_chunk(chunk_type, data):
        chunk = chunk_type + data
        return struct.pack('>I', len(data)) + chunk + struct.pack('>I', zlib.crc32(chunk) & 0xFFFFFFFF)

    signature = b'\x89PNG\r\n\x1a\n'
    ihdr_data = struct.pack('>IIBBBBB', width, height, 8, 2, 0, 0, 0)
    ihdr = make_chunk(b'IHDR', ihdr_data)

    raw_data = b''
    for y in range(height):
        raw_data += b'\x00'  # filter byte
        for x in range(width):
            raw_data += bytes(color)

    idat = make_chunk(b'IDAT', zlib.compress(raw_data))
    iend = make_chunk(b'IEND', b'')

    with open(filepath, 'wb') as f:
        f.write(signature + ihdr + idat + iend)


def add_appear_animations_only(slide, shape_ids):
    """
    Add 'Appear' entrance animations for each shape, triggered 'After Previous'.
    NO exit animations.

    Animation sequence (all in one click group, after previous):
      Object 1: Appear (afterPrevious)
      Object 2: Appear (afterPrevious)
      Object 3: Appear (afterPrevious)
      Object 4: Appear (afterPrevious)
    """
    # Build animation effect XML for each shape
    ctn_id = [1]  # mutable counter

    def next_id():
        val = ctn_id[0]
        ctn_id[0] += 1
        return val

    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # Build the animation effects for each object
    effect_pars = []
    for i, spid in enumerate(shape_ids):
        # Each appear effect is an "afterEffect" (after previous)
        par_id = next_id()
        ctn_inner_id = next_id()
        set_ctn_id = next_id()

        effect_xml = f'''<p:par xmlns:p="{ns_p}">
          <p:cTn id="{ctn_inner_id}" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" nodeType="afterEffect">
            <p:stCondLst>
              <p:cond delay="0"/>
            </p:stCondLst>
            <p:childTnLst>
              <p:set>
                <p:cBhvr>
                  <p:cTn id="{set_ctn_id}" dur="1" fill="hold">
                    <p:stCondLst>
                      <p:cond delay="0"/>
                    </p:stCondLst>
                  </p:cTn>
                  <p:tgtEl>
                    <p:spTgt spid="{spid}"/>
                  </p:tgtEl>
                  <p:attrNameLst>
                    <p:attrName>style.visibility</p:attrName>
                  </p:attrNameLst>
                </p:cBhvr>
                <p:to>
                  <p:strVal val="visible"/>
                </p:to>
              </p:set>
            </p:childTnLst>
          </p:cTn>
        </p:par>'''
        effect_pars.append(effect_xml)

    # First object uses "clickEffect" nodeType, rest use "afterEffect"
    # Actually for "After Previous" on all, the first one in a click group
    # should be nodeType="clickEffect" and the rest "afterEffect"
    # Let's fix: first shape is clickEffect (triggered by slide advance/click),
    # subsequent shapes are afterEffect (after previous finishes)
    # Re-build with correct node types
    effect_pars = []
    for i, spid in enumerate(shape_ids):
        ctn_inner_id = next_id()
        set_ctn_id = next_id()
        node_type = "clickEffect" if i == 0 else "afterEffect"

        effect_xml = f'''<p:par xmlns:p="{ns_p}">
          <p:cTn id="{ctn_inner_id}" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" nodeType="{node_type}">
            <p:stCondLst>
              <p:cond delay="0"/>
            </p:stCondLst>
            <p:childTnLst>
              <p:set>
                <p:cBhvr>
                  <p:cTn id="{set_ctn_id}" dur="1" fill="hold">
                    <p:stCondLst>
                      <p:cond delay="0"/>
                    </p:stCondLst>
                  </p:cTn>
                  <p:tgtEl>
                    <p:spTgt spid="{spid}"/>
                  </p:tgtEl>
                  <p:attrNameLst>
                    <p:attrName>style.visibility</p:attrName>
                  </p:attrNameLst>
                </p:cBhvr>
                <p:to>
                  <p:strVal val="visible"/>
                </p:to>
              </p:set>
            </p:childTnLst>
          </p:cTn>
        </p:par>'''
        effect_pars.append(effect_xml)

    effects_xml = '\n'.join(effect_pars)

    # Build the click group par
    click_group_id = next_id()
    main_seq_id = next_id()
    root_id = next_id()

    timing_xml = f'''<p:timing xmlns:p="{ns_p}" xmlns:a="{ns_a}">
  <p:tnLst>
    <p:par>
      <p:cTn id="{root_id}" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="{main_seq_id}" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                <p:par>
                  <p:cTn id="{click_group_id}" fill="hold">
                    <p:stCondLst>
                      <p:cond delay="0"/>
                    </p:stCondLst>
                    <p:childTnLst>
                      {effects_xml}
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onPrev" delay="0">
                <p:tgtEl>
                  <p:sldTgt/>
                </p:tgtEl>
              </p:cond>
            </p:prevCondLst>
            <p:nextCondLst>
              <p:cond evt="onNext" delay="0">
                <p:tgtEl>
                  <p:sldTgt/>
                </p:tgtEl>
              </p:cond>
            </p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>'''

    timing_el = etree.fromstring(timing_xml.encode())
    slide._element.append(timing_el)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slides 1-8: Realistic presentation content ---
    slide_contents = [
        ("Sequential Display Techniques", "Advanced Animation Workshop", 0),
        ("Workshop Overview", None, 5),
        ("Why Sequential Display?", None, 5),
        ("Types of Animations", None, 5),
        ("Entrance Effects", None, 5),
        ("Exit Effects", None, 5),
        ("Timing Controls", None, 5),
        ("Combining Effects", None, 5),
    ]

    body_texts = [
        None,  # slide 1 is title
        "This workshop covers advanced animation techniques in LibreOffice Impress.\n\n"
        "Topics include:\n- Entrance animations\n- Exit animations\n- Timing and sequencing\n- Professional presentation flow",
        "Sequential display helps audiences focus on one point at a time.\n\n"
        "Benefits:\n- Reduces cognitive overload\n- Controls narrative pace\n- Highlights key information\n- Creates visual hierarchy",
        "LibreOffice Impress supports several animation categories:\n\n"
        "1. Entrance - Objects appear on the slide\n"
        "2. Emphasis - Objects change while visible\n"
        "3. Exit - Objects leave the slide\n"
        "4. Motion Paths - Objects move along paths",
        "Common entrance effects:\n\n"
        "- Appear: Instant visibility\n"
        "- Fade In: Gradual opacity increase\n"
        "- Fly In: Slide from edge\n"
        "- Zoom: Scale from small to full size",
        "Common exit effects:\n\n"
        "- Disappear: Instant removal\n"
        "- Fade Out: Gradual opacity decrease\n"
        "- Fly Out: Slide to edge\n"
        "- Shrink: Scale down to nothing",
        "Timing options:\n\n"
        "- On Click: Manual trigger\n"
        "- With Previous: Simultaneous\n"
        "- After Previous: Sequential\n"
        "- Delay: Pause before animation starts",
        "Best practices for combining effects:\n\n"
        "- Pair entrance with exit for clean transitions\n"
        "- Use consistent timing across related objects\n"
        "- Add delays for readability\n"
        "- Test the full sequence before presenting",
    ]

    for i, (title, subtitle, layout_idx) in enumerate(slide_contents):
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        if slide.shapes.title:
            slide.shapes.title.text = title
        if layout_idx == 0 and subtitle:
            slide.placeholders[1].text = subtitle
        if layout_idx == 5 and body_texts[i]:
            # Add text box for body content on blank slides
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = body_texts[i]
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 9: The key slide with 4 objects and Appear animations ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    # Create a small image file for the image shape
    img_path = f'{WORKDIR}/_temp_slide9_img.png'
    create_minimal_png(img_path, 240, 180, (41, 128, 185))

    # Object 1: Image (top-left area)
    pic = slide9.shapes.add_picture(img_path, Inches(1), Inches(1), Inches(3), Inches(2.25))
    pic_id = pic.shape_id

    # Object 2: Text box (top-right area)
    txBox = slide9.shapes.add_textbox(Inches(5), Inches(1), Inches(7), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Sequential Display Demo"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    p2 = tf.add_paragraph()
    p2.text = "Each object appears one at a time, allowing the audience to focus on individual elements before the next one is revealed."
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    txBox_id = txBox.shape_id

    # Object 3: A chart-like grouped shape (bottom-left) - use a rectangle with chart title
    # We'll use a freeform/rectangle to represent a chart placeholder
    chart_shape = slide9.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4), Inches(5), Inches(3)
    )
    chart_shape.fill.solid()
    chart_shape.fill.fore_color.rgb = RGBColor(0xEC, 0xF0, 0xF1)
    chart_shape.line.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
    chart_shape.text_frame.text = "Q3 Revenue Analysis\n$2.4M Total"
    for para in chart_shape.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    chart_id = chart_shape.shape_id

    # Object 4: Geometric shape (bottom-right) - arrow/star
    arrow_shape = slide9.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(7), Inches(4.5), Inches(4), Inches(2)
    )
    arrow_shape.fill.solid()
    arrow_shape.fill.fore_color.rgb = RGBColor(0xE7, 0x4C, 0x3C)
    arrow_shape.line.color.rgb = RGBColor(0xC0, 0x39, 0x2B)
    arrow_shape.text_frame.text = "Next Steps"
    for para in arrow_shape.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(22)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    arrow_id = arrow_shape.shape_id

    # Add Appear entrance animations (After Previous) - NO exit animations
    shape_ids = [pic_id, txBox_id, chart_id, arrow_id]
    add_appear_animations_only(slide9, shape_ids)

    # --- Slide 10: Closing slide ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide10.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    p2 = tf.add_paragraph()
    p2.text = "Questions and Discussion"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Slide 9 shape IDs: pic={pic_id}, txBox={txBox_id}, chart={chart_id}, arrow={arrow_id}')
    print(f'Total slides: {len(prs.slides)}')

    # Clean up temp image
    try:
        os.remove(img_path)
    except OSError:
        pass

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
