"""
Reward Script: Competitive analysis table on slide 7
Task ID: impress_sales_057
Domain: libreoffice_impress
Scoring:
  C1 (0.15) - Table exists on slide 7, 6x5
  C2 (0.15) - Header row correct
  C3 (0.15) - Feature names in column 0
  C4 (0.15) - "Us" column all checkmarks in green
  C5 (0.25) - Competitor columns correct symbols+colors
  C6 (0.15) - Footnote text box present
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_057'

# Expected data
EXPECTED_HEADERS = ['Capability', 'Us', 'Competitor A', 'Competitor B', 'Competitor C']
EXPECTED_FEATURES = ['AI Analytics', 'Real-time Sync', 'API Access', 'Custom Workflows', 'Enterprise SSO']

# Expected symbols for rows 1-5, columns 1-4 (Us, CompA, CompB, CompC)
# True = checkmark, False = X mark
EXPECTED_GRID = [
    # Us,  CompA, CompB, CompC
    [True,  False, True,  False],  # AI Analytics
    [True,  True,  False, False],  # Real-time Sync
    [True,  True,  True,  False],  # API Access
    [True,  False, False, True],   # Custom Workflows
    [True,  True,  False, True],   # Enterprise SSO
]

GREEN_COLOR = '00AA00'
RED_COLOR = 'CC0000'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Check we have at least 7 slides
    if len(prs.slides) < 7:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 7")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[6]  # Slide 7 (0-indexed)

    # Find table shape on slide 7
    table_shape = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_shape = shape
            break

    # Component 1: Table exists on slide 7 with 6x5 dimensions (0.15 points)
    try:
        if table_shape is not None:
            table = table_shape.table
            nrows = len(table.rows)
            ncols = len(table.columns)
            if nrows == 6 and ncols == 5:
                print(f"PASS: Component 1 - Table found on slide 7, dimensions 6x5 (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 1 - Table dimensions {nrows}x{ncols}, expected 6x5")
        else:
            print("FAIL: Component 1 - No table found on slide 7")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    if table_shape is None:
        # No table means all remaining checks fail
        print("REWARD: 0.0")
        return 0.0

    table = table_shape.table

    # Component 2: Header row matches expected values (0.15 points)
    try:
        headers = [table.cell(0, c).text.strip() for c in range(5)]
        if headers == EXPECTED_HEADERS:
            print(f"PASS: Component 2 - Header row correct: {headers} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 - Headers: {headers}, expected: {EXPECTED_HEADERS}")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Feature names in column 0, rows 1-5 (0.15 points)
    try:
        features = [table.cell(r, 0).text.strip() for r in range(1, 6)]
        if features == EXPECTED_FEATURES:
            print(f"PASS: Component 3 - Feature names correct (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 3 - Features: {features}, expected: {EXPECTED_FEATURES}")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: "Us" column (col 1) has all checkmarks in green (0.15 points)
    try:
        us_pass = True
        for r in range(1, 6):
            cell = table.cell(r, 1)
            cell_text = cell.text.strip()
            # Check for checkmark symbol
            if '\u2713' not in cell_text:
                print(f"FAIL: Component 4 - Row {r} 'Us' column: expected checkmark, found {repr(cell_text)}")
                us_pass = False
                break
            # Check color is green
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    if '\u2713' in run.text:
                        try:
                            clr = str(run.font.color.rgb) if run.font.color.type is not None else 'None'
                        except:
                            clr = 'None'
                        if clr != GREEN_COLOR:
                            print(f"FAIL: Component 4 - Row {r} 'Us' checkmark color: {clr}, expected {GREEN_COLOR}")
                            us_pass = False
                            break
                if not us_pass:
                    break
            if not us_pass:
                break

        if us_pass:
            print(f"PASS: Component 4 - 'Us' column all green checkmarks (0.15 pts)")
            total_score += 0.15
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Competitor columns correct symbols and colors (0.25 points)
    try:
        comp_correct = 0
        comp_total = 15  # 5 rows x 3 competitor columns
        for r in range(5):
            for c_offset in range(3):  # columns 2, 3, 4
                c = c_offset + 2
                cell = table.cell(r + 1, c)
                cell_text = cell.text.strip()
                expected_check = EXPECTED_GRID[r][c_offset + 1]

                if expected_check:
                    expected_sym = '\u2713'
                    expected_color = GREEN_COLOR
                else:
                    expected_sym = '\u2717'
                    expected_color = RED_COLOR

                sym_ok = expected_sym in cell_text
                color_ok = False
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if expected_sym in run.text:
                            try:
                                clr = str(run.font.color.rgb) if run.font.color.type is not None else 'None'
                            except:
                                clr = 'None'
                            if clr == expected_color:
                                color_ok = True

                if sym_ok and color_ok:
                    comp_correct += 1
                else:
                    feat = EXPECTED_FEATURES[r]
                    col_name = EXPECTED_HEADERS[c]
                    print(f"  DETAIL: [{feat}, {col_name}] sym_ok={sym_ok} color_ok={color_ok} text={repr(cell_text)}")

        if comp_correct == comp_total:
            print(f"PASS: Component 5 - All competitor cells correct ({comp_correct}/{comp_total}) (0.25 pts)")
            total_score += 0.25
        elif comp_correct >= 10:
            partial = round(0.25 * comp_correct / comp_total, 2)
            print(f"PARTIAL: Component 5 - {comp_correct}/{comp_total} competitor cells correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 - Only {comp_correct}/{comp_total} competitor cells correct")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    # Component 6: Footnote text box present with correct text (0.15 points)
    try:
        footnote_found = False
        expected_footnote = 'Source: Independent analysis, March 2025'
        for shape in slide.shapes:
            if shape.has_text_frame and shape.shape_type != MSO_SHAPE_TYPE.TABLE:
                txt = shape.text_frame.text.strip()
                if 'Source:' in txt and 'March 2025' in txt:
                    footnote_found = True
                    if expected_footnote.lower() in txt.lower():
                        print(f"PASS: Component 6 - Footnote found: {repr(txt)} (0.15 pts)")
                        total_score += 0.15
                    else:
                        # Partial: has source reference but not exact
                        print(f"PARTIAL: Component 6 - Footnote close but not exact: {repr(txt)} (0.10 pts)")
                        total_score += 0.10
                    break

        if not footnote_found:
            # Also check for partial footnote
            for shape in slide.shapes:
                if shape.has_text_frame and shape.shape_type != MSO_SHAPE_TYPE.TABLE:
                    txt = shape.text_frame.text.strip()
                    if 'independent analysis' in txt.lower():
                        footnote_found = True
                        print(f"PARTIAL: Component 6 - Partial footnote found: {repr(txt)} (0.05 pts)")
                        total_score += 0.05
                        break

            if not footnote_found:
                print(f"FAIL: Component 6 - No footnote text box found on slide 7")
    except Exception as e:
        print(f"ERROR: Component 6 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
