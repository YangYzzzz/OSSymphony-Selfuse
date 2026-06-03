"""
Initial Setup: Competitive Analysis Pitch Deck with empty slide 7
Task ID: impress_sales_057
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_057'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
    return slide


def add_title_only_slide(prs, title_text):
    """Add a slide with title only and no body content."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8.5), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    add_title_slide(prs, "CompAnalysis Platform", "Strategic Market Positioning & Competitive Landscape\nQ1 2025 Board Review")

    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "Market share grew 18% YoY reaching $47.2M ARR",
        "Enterprise segment expanded with 23 new logos in Q4",
        "Product NPS improved to 72 from 64 in prior quarter",
        "Churn rate decreased to 3.2% from 4.8% annually",
    ])

    # Slide 3: Market Overview
    add_content_slide(prs, "Market Overview", [
        "Total addressable market estimated at $12.8B by 2026",
        "SaaS analytics segment growing at 24% CAGR",
        "Key verticals: Financial Services, Healthcare, Retail",
        "Emerging demand for AI-driven predictive analytics",
        "Regulatory compliance driving enterprise adoption",
    ])

    # Slide 4: Our Product Strengths
    add_content_slide(prs, "Our Product Strengths", [
        "Proprietary AI engine with 99.7% prediction accuracy",
        "Sub-100ms real-time data synchronization",
        "RESTful API with 200+ endpoints and SDKs in 8 languages",
        "Custom workflow builder used by 89% of enterprise clients",
        "SOC 2 Type II and ISO 27001 certified",
    ])

    # Slide 5: Competitor Landscape
    add_content_slide(prs, "Competitor Landscape", [
        "Competitor A: Legacy analytics provider, strong in finance vertical",
        "Competitor B: Cloud-native startup, aggressive pricing model",
        "Competitor C: Enterprise suite vendor, broad but shallow features",
        "Key differentiator: Our end-to-end AI integration",
    ])

    # Slide 6: Pricing Strategy
    add_content_slide(prs, "Pricing Strategy", [
        "Starter: $49/user/month - Core analytics and dashboards",
        "Professional: $129/user/month - AI features and API access",
        "Enterprise: Custom pricing - Full platform with SSO and SLA",
        "Volume discounts available for 500+ seat deployments",
    ])

    # Slide 7: Feature Comparison - TITLE ONLY, NO TABLE (task target)
    add_title_only_slide(prs, "Feature Comparison")

    # Slide 8: Customer Testimonials
    add_content_slide(prs, "Customer Testimonials", [
        '"Reduced our reporting cycle from 2 weeks to 2 hours" - Sarah Kim, VP Analytics, Meridian Health',
        '"The AI recommendations increased our conversion rate by 34%" - David Park, CTO, NovaTech',
        '"Best-in-class API documentation and developer experience" - Lisa Torres, Lead Engineer, FinCore',
    ])

    # Slide 9: Growth Roadmap
    add_content_slide(prs, "Growth Roadmap 2025-2026", [
        "Q2 2025: Launch generative AI report builder",
        "Q3 2025: Expand to APAC with localized deployment",
        "Q4 2025: Achieve FedRAMP authorization for gov sector",
        "Q1 2026: Release mobile-first analytics dashboard",
        "Q2 2026: IPO readiness milestone",
    ])

    # Slide 10: Next Steps
    add_content_slide(prs, "Next Steps & Action Items", [
        "Finalize Series C funding round by end of April",
        "Onboard 15 new enterprise accounts in Q2 pipeline",
        "Accelerate Competitor A displacement campaigns",
        "Hire 20 additional engineers for AI team expansion",
        "Schedule quarterly board review for July 15, 2025",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
