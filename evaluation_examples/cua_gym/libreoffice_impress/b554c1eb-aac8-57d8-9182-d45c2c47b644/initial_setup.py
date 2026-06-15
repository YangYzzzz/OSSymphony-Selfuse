"""
Initial Setup: Create a 12-slide presentation for handout export task
Task ID: impstruct_043
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impstruct_043'
OUTPUT = f'{WORKDIR}/handout_deck.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.paragraphs[0].text = point
        else:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
    return slide


def create_initial():
    prs = Presentation()

    # Slide 1: Title slide
    add_title_slide(prs, "Q1 2025 Business Review", "Meridian Analytics Corporation")

    # Slide 2: Agenda
    add_content_slide(prs, "Meeting Agenda", [
        "Revenue performance overview",
        "Customer acquisition metrics",
        "Product development milestones",
        "Regional expansion update",
        "Q2 strategic priorities",
    ])

    # Slide 3: Revenue Overview
    add_content_slide(prs, "Revenue Overview", [
        "Total revenue: $14.7M (+18% YoY)",
        "Recurring revenue: $11.2M (76% of total)",
        "Average deal size increased to $47,500",
        "Enterprise segment grew 32%",
    ])

    # Slide 4: Customer Metrics
    add_content_slide(prs, "Customer Acquisition & Retention", [
        "New customers: 238 (up from 194 in Q4)",
        "Customer churn rate: 2.1% (industry avg: 5.8%)",
        "Net Promoter Score: 72",
        "Customer lifetime value: $185,000",
    ])

    # Slide 5: Product Development
    add_content_slide(prs, "Product Development Milestones", [
        "Launched Analytics Dashboard v3.2",
        "Mobile app beta released to 500 users",
        "API response time reduced by 40%",
        "New integrations: Salesforce, HubSpot, Jira",
    ])

    # Slide 6: Engineering Team
    add_content_slide(prs, "Engineering Highlights", [
        "Team grew from 45 to 58 engineers",
        "Deployed 147 releases with 99.95% uptime",
        "Technical debt reduced by 22%",
        "Migrated 3 core services to Kubernetes",
    ])

    # Slide 7: Marketing Performance
    add_content_slide(prs, "Marketing Performance", [
        "Marketing qualified leads: 1,847",
        "Cost per acquisition: $312 (down 15%)",
        "Website traffic: 892K unique visitors",
        "Webinar attendance averaged 340 per session",
    ])

    # Slide 8: Regional Expansion
    add_content_slide(prs, "Regional Expansion Update", [
        "EMEA office opened in Frankfurt, Germany",
        "APAC revenue grew 45% quarter-over-quarter",
        "Hired regional sales directors for 4 markets",
        "Localization complete for 8 languages",
    ])

    # Slide 9: Financial Summary
    add_content_slide(prs, "Financial Summary", [
        "Gross margin: 78.3% (target: 75%)",
        "Operating expenses: $9.8M",
        "EBITDA: $3.2M (positive for 3rd consecutive quarter)",
        "Cash runway: 24+ months",
    ])

    # Slide 10: Key Partnerships
    add_content_slide(prs, "Strategic Partnerships", [
        "Signed enterprise deal with Siemens ($2.1M ARR)",
        "Technology partnership with AWS Advanced tier",
        "Channel partnership with Deloitte Consulting",
        "Academic collaboration with MIT Media Lab",
    ])

    # Slide 11: Q2 Priorities
    add_content_slide(prs, "Q2 2025 Strategic Priorities", [
        "Launch enterprise self-service portal",
        "Expand APAC sales team by 12 headcount",
        "Achieve SOC 2 Type II certification",
        "Release AI-powered anomaly detection feature",
        "Target $16.5M quarterly revenue",
    ])

    # Slide 12: Closing
    add_title_slide(prs, "Thank You", "Questions & Discussion\ncontact@meridiananalytics.com")

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
