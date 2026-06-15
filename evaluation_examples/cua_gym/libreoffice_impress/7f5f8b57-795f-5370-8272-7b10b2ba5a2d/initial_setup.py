"""
Initial Setup: Create a presentation with a bar chart on slide 2
Task ID: impress_tct_055
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_055'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Quarterly Revenue Analysis"
    slide1.placeholders[1].text = "Prepared by Analytics Division — Q1 2025"

    # --- Slide 2: Bar Chart ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Add a title text box
    txBox = slide2.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Revenue by Product Line"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Chart data
    chart_data = CategoryChartData()
    chart_data.categories = ['Enterprise', 'SMB', 'Consumer', 'Government', 'Education']
    chart_data.add_series('Q4 2024', (285000, 142000, 198000, 95000, 67000))
    chart_data.add_series('Q1 2025', (312000, 158000, 215000, 102000, 73000))

    chart_shape = slide2.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.0), Inches(1.5), Inches(10), Inches(5.2),
        chart_data
    )
    chart = chart_shape.chart

    # Style the chart — white background (default), no plot area border
    # Set chart area fill to white explicitly
    chart_element = chart._chartSpace
    chart_area = chart_element.find(qn('c:chart'))

    # Ensure chart background is explicitly white (not light gray)
    # The default is no fill, which appears white. We leave it as default.

    # Style series colors
    plot = chart.plots[0]
    series0 = plot.series[0]
    series0.format.fill.solid()
    series0.format.fill.fore_color.rgb = RGBColor(0x5B, 0x9B, 0xD5)
    series1 = plot.series[1]
    series1.format.fill.solid()
    series1.format.fill.fore_color.rgb = RGBColor(0x70, 0xAD, 0x47)

    # Chart has legend
    chart.has_legend = True
    chart.legend.include_in_layout = False

    # --- Slide 3: Summary Table ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    txBox3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(8), Inches(0.8))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Revenue Growth Summary"
    p3.alignment = PP_ALIGN.LEFT
    run3 = p3.runs[0]
    run3.font.name = "Calibri"
    run3.font.size = Pt(28)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    rows, cols = 6, 4
    table_shape = slide3.shapes.add_table(rows, cols, Inches(1.0), Inches(1.5), Inches(10), Inches(4.0))
    table = table_shape.table
    headers = ['Product Line', 'Q4 2024', 'Q1 2025', 'Growth %']
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)

    data_rows = [
        ['Enterprise', '$285,000', '$312,000', '+9.5%'],
        ['SMB', '$142,000', '$158,000', '+11.3%'],
        ['Consumer', '$198,000', '$215,000', '+8.6%'],
        ['Government', '$95,000', '$102,000', '+7.4%'],
        ['Education', '$67,000', '$73,000', '+9.0%'],
    ]
    for r, row_data in enumerate(data_rows, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 4: Key Takeaways ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    txBox4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(8), Inches(0.8))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Key Takeaways"
    p4.alignment = PP_ALIGN.LEFT
    run4 = p4.runs[0]
    run4.font.name = "Calibri"
    run4.font.size = Pt(28)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    bullets = [
        "Enterprise segment leads with $312K in Q1, driven by new SaaS contracts",
        "SMB shows strongest growth at 11.3%, reflecting improved onboarding",
        "Consumer revenue reached $215K, up from seasonal holiday momentum",
        "Government and Education segments remain stable with steady expansion",
        "Overall portfolio grew 9.2% quarter-over-quarter across all segments",
    ]
    bullet_box = slide4.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(10), Inches(5.0))
    btf = bullet_box.text_frame
    btf.word_wrap = True
    for i, bullet_text in enumerate(bullets):
        if i == 0:
            p = btf.paragraphs[0]
        else:
            p = btf.add_paragraph()
        p.text = bullet_text
        p.level = 0
        p.space_after = Pt(12)
        for run in p.runs:
            run.font.size = Pt(16)
            run.font.name = "Calibri"
            run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
