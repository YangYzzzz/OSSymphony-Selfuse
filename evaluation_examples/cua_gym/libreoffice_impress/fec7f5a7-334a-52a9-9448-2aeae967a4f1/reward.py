"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 71 I’ve got three text boxes that need tidying up. How do I get LibreOffice Impress to line all their left edges up perfectly and then spread them out evenly from top to bottom?
Generated: 2025-09-11 00:27:56
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
import statistics
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def verify_impress_alignment(file_path: str) -> float:
    """Verify that on slide 71 the three text-boxes
    1) share the same left edge (are vertically aligned)
    2) are evenly distributed from top to bottom.

    Returns a progressive score between 0.0 and 1.0 and prints
    detailed diagnostics for each verification step.
    """

    print(f"Loading presentation: {file_path}\n")

    # ---------- prerequisite checks (NO POINTS AWARDED) ----------
    if not os.path.exists(file_path):
        print("✗ File not found – task failed")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Could not open presentation: {e}")
        return 0.0

    # Ensure we have at least 71 slides (index 70)
    if len(prs.slides) <= 70:
        print(f"✗ Presentation has only {len(prs.slides)} slides – slide 71 missing")
        return 0.0

    slide_71 = prs.slides[70]

    # ---------- locate text-boxes on slide 71 ----------
    text_boxes = [sh for sh in slide_71.shapes
                  if sh.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and sh.has_text_frame]

    print(f"Found {len(text_boxes)} text boxes on slide 71")

    total_score = 0.0  # progressive scoring starts here

    # --- Requirement 1: at least three text boxes (0.2 points) ---
    if len(text_boxes) >= 3:
        total_score += 0.2
        print("✓ At least three text boxes present (0.2 points)")
    else:
        print("✗ Fewer than three text boxes – cannot fulfil task requirements")
        print(f"REWARD: {total_score}")
        return total_score  # early exit – no further points possible

    # We'll analyse all the text boxes present (could be >3)

    # --- Requirement 2: left-edge alignment (0.4 points) ---
    left_positions = [tb.left for tb in text_boxes]
    median_left = statistics.median(left_positions)
    # Allow small tolerance (5,000 EMU ≈ 0.07 inches)
    alignment_tolerance = 5000
    aligned = all(abs(l - median_left) <= alignment_tolerance for l in left_positions)

    if aligned:
        total_score += 0.4
        print("✓ All text boxes share the same left edge (0.4 points)")
    else:
        print("✗ Text boxes are NOT aligned on the left edge (0 points)")

    # --- Requirement 3: even vertical distribution (0.4 points) ---
    # Sort by top coordinate to analyse spacing
    tops_sorted = sorted([tb.top for tb in text_boxes])
    if len(tops_sorted) >= 3:
        gaps = [tops_sorted[i + 1] - tops_sorted[i] for i in range(len(tops_sorted) - 1)]
        avg_gap = sum(gaps) / len(gaps)
        even_tolerance = 5000  # same 0.07" tolerance for spacing equality
        even_spacing = all(abs(g - avg_gap) <= even_tolerance for g in gaps)

        if even_spacing:
            total_score += 0.4
            print("✓ Text boxes are evenly distributed vertically (0.4 points)")
        else:
            print("✗ Text boxes are NOT evenly distributed (0 points)")
    else:
        # Should not occur because we already checked len >=3, but safe guard
        print("✗ Not enough text boxes to evaluate distribution (0 points)")

    final_score = min(total_score, 1.0)
    print(f"\nFinal score: {final_score}")
    return final_score


if __name__ == "__main__":
    FILE_PATH = "/home/user/on_slide_71_ive_got_three_text_boxes_that_need_tidying_up_how_do_i_get_libreoffice_impress_to_line_a_golden.pptx"
    reward_value = verify_impress_alignment(FILE_PATH)
    print(f"REWARD: {reward_value}")

