"""
Initial Setup: Duplicate slide 2 and place the copy right after it as slide 3.
Task ID: osworld_impress_slide_duplication_reorder_001
Domain: libreoffice_impress

Creates a 5-slide business deck. Slide 2 contains a 'Strategy Overview'
layout with a title and bullet list. The agent must duplicate slide 2 and
insert the copy as slide 3, resulting in a 6-slide deck.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_slide_duplication_reorder_001'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_bullet_text(tf, title_text, bullets):
    """Helper: set title placeholder text and populate bullets in a content frame."""
    tf.text = bullets[0]
    for bullet in bullets[1:]:
        para = tf.add_paragraph()
        para.text = bullet
        para.level = 1


def create_initial():
    prs = Presentation()
    # Standard widescreen layout (default template)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
    slide1.shapes.title.text = "Q3 Business Review"
    slide1.placeholders[1].text = "Nexora Technologies — September 2025"

    # ---- Slide 2: Strategy Overview (Title + Content) ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide2.shapes.title.text = "Strategy Overview"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Expand into Southeast Asian markets by Q4 2025"
    paras = [
        "Launch three new enterprise product lines",
        "Achieve 25% year-over-year revenue growth",
        "Reduce operational costs by 12% through automation",
        "Strengthen strategic partnerships with regional distributors",
    ]
    for bullet in paras:
        p = tf2.add_paragraph()
        p.text = bullet
        p.level = 1

    # ---- Slide 3: Market Analysis ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Market Analysis"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Total addressable market: $2.4B by 2026"
    market_bullets = [
        "Southeast Asia CAGR projected at 18.3%",
        "Key segments: Fintech, Healthcare IT, Logistics",
        "Competitive landscape: 5 major players identified",
        "Customer acquisition cost down 9% YoY",
    ]
    for bullet in market_bullets:
        p = tf3.add_paragraph()
        p.text = bullet
        p.level = 1

    # ---- Slide 4: Financial Performance ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Financial Performance"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Q3 2025 Total Revenue: $47.8M (+22% YoY)"
    fin_bullets = [
        "Gross margin improved to 64% from 61% last year",
        "EBITDA: $12.1M, exceeding forecast by 8%",
        "R&D investment: $6.5M (up 15% from Q3 2024)",
        "Cash reserves: $31.2M — sufficient for planned expansion",
    ]
    for bullet in fin_bullets:
        p = tf4.add_paragraph()
        p.text = bullet
        p.level = 1

    # ---- Slide 5: Next Steps ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Next Steps"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Immediate priorities for Q4 2025"
    next_bullets = [
        "Finalize Singapore office lease by October 15",
        "Hire 40 local enterprise sales representatives",
        "Complete product localization for Thai and Vietnamese markets",
        "Board approval for $8M regional marketing budget",
    ]
    for bullet in next_bullets:
        p = tf5.add_paragraph()
        p.text = bullet
        p.level = 1

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Slide count: {len(prs.slides)} (expected 5)')

    # GUI-ready startup: open the file in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
