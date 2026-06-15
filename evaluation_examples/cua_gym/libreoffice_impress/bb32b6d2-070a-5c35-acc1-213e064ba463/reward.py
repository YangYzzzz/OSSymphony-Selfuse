"""
Reward Script: Duplicate slide 2 and place the copy right after it as slide 3.
Task ID: osworld_impress_slide_duplication_reorder_001
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): Presentation has exactly 6 slides (was 5 initially)
  Component 2 (0.4): Slide 3 is an exact duplicate of slide 2 (same title + same bullet content)
  Component 3 (0.3): Original slides 3-5 are now at positions 4-6 (ordering preserved)
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_slide_duplication_reorder_001'

# Expected slide titles in order after task completion
EXPECTED_SLIDE_TITLES = [
    'Q3 Business Review',
    'Strategy Overview',
    'Strategy Overview',   # the duplicate
    'Market Analysis',
    'Financial Performance',
    'Next Steps',
]

# Expected content of slide 2 (and therefore the duplicate slide 3)
SLIDE2_TITLE = 'Strategy Overview'
SLIDE2_BULLETS = [
    'Expand into Southeast Asian markets by Q4 2025',
    'Launch three new enterprise product lines',
    'Achieve 25% year-over-year revenue growth',
    'Reduce operational costs by 12% through automation',
    'Strengthen strategic partnerships with regional distributors',
]


def get_slide_title(slide):
    """Return the title text of a slide, or empty string if not found."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name.startswith('Title'):
            return shape.text_frame.text.strip()
    return ''


def get_slide_bullet_texts(slide):
    """Return list of non-empty paragraph texts from non-title text frames."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame and not shape.name.startswith('Title'):
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return texts


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Presentation has exactly 6 slides (0.3 points)
    # The task adds one slide (duplicate), so 5 → 6.
    try:
        if num_slides == 6:
            print(f"PASS: Component 1 — Slide count is 6 (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 — Expected 6 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 3 is an exact duplicate of slide 2 (0.4 points)
    # Slide 3 must have the same title and the same bullet content as slide 2.
    try:
        if num_slides >= 3:
            slide2 = prs.slides[1]
            slide3 = prs.slides[2]

            slide2_title = get_slide_title(slide2)
            slide3_title = get_slide_title(slide3)
            slide2_bullets = get_slide_bullet_texts(slide2)
            slide3_bullets = get_slide_bullet_texts(slide3)

            title_match = (slide3_title == SLIDE2_TITLE)
            bullets_match = (slide3_bullets == slide2_bullets) and (len(slide3_bullets) > 0)

            if title_match and bullets_match:
                print(f"PASS: Component 2 — Slide 3 is a duplicate of slide 2 "
                      f"(title='{slide3_title}', {len(slide3_bullets)} bullets match) (0.4 pts)")
                total_score += 0.4
            else:
                if not title_match:
                    print(f"FAIL: Component 2 — Slide 3 title mismatch: "
                          f"expected '{SLIDE2_TITLE}', found '{slide3_title}'")
                if not bullets_match:
                    print(f"FAIL: Component 2 — Slide 3 bullet content mismatch: "
                          f"expected {SLIDE2_BULLETS}, found {slide3_bullets}")
        else:
            print(f"FAIL: Component 2 — Not enough slides ({num_slides}) to check slide 3")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Original slides 3-5 are preserved as slides 4-6 (0.3 points)
    # After inserting the duplicate at position 3, the original slide 3 ('Market Analysis'),
    # slide 4 ('Financial Performance'), and slide 5 ('Next Steps') should be at positions 4, 5, 6.
    try:
        if num_slides >= 6:
            expected_tail = EXPECTED_SLIDE_TITLES[3:]   # ['Market Analysis', 'Financial Performance', 'Next Steps']
            actual_tail = [get_slide_title(prs.slides[i]) for i in range(3, 6)]

            if actual_tail == expected_tail:
                print(f"PASS: Component 3 — Slides 4-6 are the original slides 3-5 "
                      f"({actual_tail}) (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 3 — Expected slides 4-6 to be "
                      f"{expected_tail}, found {actual_tail}")
        else:
            print(f"FAIL: Component 3 — Not enough slides ({num_slides}) to check slides 4-6")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the VM environment
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
