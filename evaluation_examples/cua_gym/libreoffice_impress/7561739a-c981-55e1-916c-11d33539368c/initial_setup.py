"""
Initial Setup: 6-slide marketing deck with all titles in black
Task ID: osworld_impress_title_selective_formatting_005
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_selective_formatting_005'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # Slide dimensions (default widescreen 10x7.5 inches)
    # Use default slide layouts

    # --- Slide data: 6 marketing slides ---
    slides_data = [
        {
            "title": "Q1 Marketing Strategy Overview",
            "content": (
                "Our Q1 strategy focuses on expanding brand reach, "
                "increasing digital engagement, and driving revenue growth "
                "through targeted campaigns across all major channels."
            ),
        },
        {
            "title": "Target Audience Analysis",
            "content": (
                "Primary demographics: 25-44 professionals in tech and finance sectors. "
                "Key segments include urban millennials and Gen X decision-makers "
                "with annual household incomes above $75,000."
            ),
        },
        {
            "title": "Digital Campaign Initiatives",
            "content": (
                "Launching three flagship campaigns: SpringForward social media push, "
                "email nurture sequences for existing leads, and paid search expansion "
                "targeting competitor keywords in top-performing regions."
            ),
        },
        {
            "title": "Budget Allocation & ROI Targets",
            "content": (
                "Total Q1 budget: $420,000. Allocation: 40% digital ads, "
                "25% content creation, 20% events and sponsorships, "
                "15% analytics and tooling. Target ROI: 3.2x by end of quarter."
            ),
        },
        {
            "title": "Partnership & Influencer Program",
            "content": (
                "Partnering with 12 industry influencers and 3 strategic brand allies. "
                "Focus on authentic storytelling and co-branded content that resonates "
                "with our core audience values: innovation, sustainability, and impact."
            ),
        },
        {
            "title": "Performance Metrics & Review Schedule",
            "content": (
                "KPIs tracked weekly: impressions, CTR, conversions, CAC, and LTV. "
                "Monthly review sessions with department heads. "
                "Final Q1 report due April 15 for executive presentation."
            ),
        },
    ]

    for slide_data in slides_data:
        # Use Title and Content layout (index 1)
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # Set title — all titles in BLACK (default, no explicit color override)
        title_shape = slide.shapes.title
        title_shape.text = slide_data["title"]

        # Explicitly set title color to black to be unambiguous
        for para in title_shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
                run.font.size = Pt(36)
                run.font.bold = True

        # Set content
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = slide_data["content"]
        for para in content_placeholder.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(20)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open the file in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
