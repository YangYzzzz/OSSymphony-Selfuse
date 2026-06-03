"""
Reward Script: Change title color to blue on every other slide (slides 1, 3, 5)
Task ID: osworld_impress_title_selective_formatting_005
Domain: libreoffice_impress
Scoring:
  Component 1: Slides 1, 3, 5 title color changed to blue (#0000FF) — 0.6 pts total (0.2 pts each)
  Component 2: Slides 2, 4, 6 title color remains black (#000000) — 0.4 pts total (gated on component 1 > 0)
  Total: 1.0
"""

import os
from pptx import Presentation
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_selective_formatting_005'

BLUE = RGBColor(0x00, 0x00, 0xFF)   # #0000FF
BLACK = RGBColor(0x00, 0x00, 0x00)  # #000000


def get_title_shape(slide):
    """Return the title shape for a slide, or None if not found."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name.startswith('Title'):
            return shape
    return None


def get_title_color(slide):
    """
    Return the RGB color string of all non-empty runs in the title shape.
    If runs have differing colors, return a list of color strings.
    If no color is explicitly set (type is None), return None.
    Returns a single string if all runs share the same color.
    """
    title_shape = get_title_shape(slide)
    if title_shape is None:
        return None

    colors = []
    for para in title_shape.text_frame.paragraphs:
        for run in para.runs:
            if not (run.text or '').strip():
                continue
            try:
                if run.font.color.type is not None:
                    colors.append(str(run.font.color.rgb).upper())
                else:
                    colors.append(None)  # inherited / no explicit color
            except Exception:
                colors.append(None)

    if not colors:
        return None
    # If all same, return single value
    if len(set(c for c in colors)) == 1:
        return colors[0]
    return colors  # mixed


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have exactly 6 slides
    if len(prs.slides) != 6:
        print(f"CRITICAL: Expected 6 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # ------------------------------------------------------------------
    # Component 1: Slides 1, 3, 5 titles must be blue (#0000FF)  — 0.2 pts each
    # These are the slides that should have been changed by the task.
    # On initial_env all titles are black, so this FAILS on initial.
    # ------------------------------------------------------------------
    blue_target_slides = [1, 3, 5]   # 1-indexed
    for slide_num in blue_target_slides:
        slide = prs.slides[slide_num - 1]
        try:
            color = get_title_color(slide)
            if color == '0000FF':
                print(f"PASS: Slide {slide_num} title is blue (#0000FF) (+0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Slide {slide_num} title expected blue (#0000FF), found: {color!r}")
        except Exception as e:
            print(f"ERROR: Could not check title color on slide {slide_num}: {e}")

    # ------------------------------------------------------------------
    # Component 2: Slides 2, 4, 6 titles must remain black (#000000) — 0.4 pts total
    # These slides should NOT have been modified. Checking them ensures the task
    # was done selectively. On initial_env all titles are black → this check would
    # pass in isolation. However this is a COMPOUND check: points are only awarded
    # when at least one blue-slide check already passed (i.e., the task was at least
    # partially attempted). We gate on total_score > 0 to avoid giving 0.4 pts on
    # the initial_env where nothing was done yet.
    # ------------------------------------------------------------------
    black_target_slides = [2, 4, 6]   # 1-indexed; should stay black
    black_correct_count = 0
    for slide_num in black_target_slides:
        slide = prs.slides[slide_num - 1]
        try:
            color = get_title_color(slide)
            # Accept both explicit black (000000) and inherited (None = defaults to black)
            if color == '000000' or color is None:
                print(f"PASS: Slide {slide_num} title remains black ({color!r}) (pending gate)")
                black_correct_count += 1
            else:
                print(f"FAIL: Slide {slide_num} title should remain black, found: {color!r}")
        except Exception as e:
            print(f"ERROR: Could not check title color on slide {slide_num}: {e}")

    # Gate: only award black-slide points when at least one blue slide was changed.
    # This prevents awarding points on the initial_env where nothing was done yet.
    if total_score > 0.0 and black_correct_count == 3:
        total_score += 0.4
        print(f"PASS: All 3 even slides remain black — selective formatting verified (gate open — +0.4 pts)")
    elif black_correct_count == 3:
        print(f"INFO: Even slides are black, but no blue slides changed yet — gate closed (0.0 pts)")
    else:
        print(f"FAIL: {3 - black_correct_count} even slide(s) have unexpected color changes — no component 2 pts")

    final_score = round(min(total_score, 1.0), 4)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
