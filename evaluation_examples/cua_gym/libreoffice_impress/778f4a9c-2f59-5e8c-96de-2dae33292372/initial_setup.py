"""
Initial Setup: Sales pitch deck with 5 slides — pre-task state.
Task ID: osworld_impress_multi_op_combined_004
Domain: libreoffice_impress

Initial state:
  - 5-slide sales pitch deck
  - Slide 2 title "Our Solution" is plain black (NOT bold, NOT blue)
  - Slide 3 has a title and body content but NO additional textboxes
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_multi_op_combined_004'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Accelerate Your Growth"
    slide1.placeholders[1].text = "A Strategic Sales Pitch by NovaTech Solutions"

    # ---- Slide 2: Our Solution ----
    # Title must be plain black — NOT bold, NOT blue
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    title2.text = "Our Solution"
    # Explicitly set plain black, not bold
    for para in title2.text_frame.paragraphs:
        for run in para.runs:
            run.font.bold = False
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    body2 = slide2.placeholders[1]
    tf2 = body2.text_frame
    tf2.text = "NovaTech's platform addresses your core pain points with:"
    p2a = tf2.add_paragraph()
    p2a.text = "Automated workflow integration across 50+ enterprise tools"
    p2a.level = 1
    p2b = tf2.add_paragraph()
    p2b.text = "Real-time analytics dashboard with AI-driven recommendations"
    p2b.level = 1
    p2c = tf2.add_paragraph()
    p2c.text = "Zero-downtime deployment and 99.9% uptime SLA"
    p2c.level = 1
    p2d = tf2.add_paragraph()
    p2d.text = "Dedicated customer success team with 24/7 support"
    p2d.level = 1

    # ---- Slide 3: Market Opportunity ----
    # Title + body content only, NO additional textboxes
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = "Market Opportunity"

    body3 = slide3.placeholders[1]
    tf3 = body3.text_frame
    tf3.text = "The total addressable market is growing at 23% CAGR:"
    p3a = tf3.add_paragraph()
    p3a.text = "Global enterprise software market: $650B by 2026"
    p3a.level = 1
    p3b = tf3.add_paragraph()
    p3b.text = "Workflow automation segment: $78B and expanding rapidly"
    p3b.level = 1
    p3c = tf3.add_paragraph()
    p3c.text = "Early movers capturing 3-5x greater ROI vs. late adopters"
    p3c.level = 1
    p3d = tf3.add_paragraph()
    p3d.text = "Key verticals: Finance, Healthcare, Retail, Manufacturing"
    p3d.level = 1

    # ---- Slide 4: Case Studies ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    title4.text = "Customer Success Stories"

    body4 = slide4.placeholders[1]
    tf4 = body4.text_frame
    tf4.text = "Proven results across diverse industries:"
    p4a = tf4.add_paragraph()
    p4a.text = "Meridian Financial: 40% reduction in manual processing time"
    p4a.level = 1
    p4b = tf4.add_paragraph()
    p4b.text = "HealthBridge Clinic Network: $2.3M saved in operational costs"
    p4b.level = 1
    p4c = tf4.add_paragraph()
    p4c.text = "RetailMax Group: 18% revenue uplift from automated campaigns"
    p4c.level = 1
    p4d = tf4.add_paragraph()
    p4d.text = "SteelCore Manufacturing: 99.97% uptime with zero unplanned outages"
    p4d.level = 1

    # ---- Slide 5: Call to Action ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[0])
    slide5.shapes.title.text = "Next Steps"
    slide5.placeholders[1].text = (
        "Schedule a personalized demo today.\n"
        "Contact: sales@novatech.io  |  +1 (800) 555-0192"
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
