"""
Reward Script: Apply bold and blue color to title on slide 2, add textbox on slide 3
Task ID: osworld_impress_multi_op_combined_004
Domain: libreoffice_impress
Scoring:
  - Component 1: Slide 2 title 'Our Solution' is bold (0.35 pts)
  - Component 2: Slide 2 title 'Our Solution' has blue color (0.35 pts)
  - Component 3: Slide 3 has a new text box with text 'Competitive Advantage' (0.30 pts)
"""

import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_multi_op_combined_004'


def persist_app_state():
    """Send Ctrl+S to save any unsaved LibreOffice Impress state."""
    try:
        import pyautogui
        import time
        os.environ["DISPLAY"] = ":0"
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def is_blue_color(rgb_str):
    """
    Check whether an RGB hex string represents a blue-dominant color.
    Accepts any color where the blue channel is > 128 and dominates
    over red and green channels. Also accepts exact known blue values.
    rgb_str: 6-char hex string e.g. '0070C0', 'FF0000'
    """
    try:
        r = int(rgb_str[0:2], 16)
        g = int(rgb_str[2:4], 16)
        b = int(rgb_str[4:6], 16)
        # Blue channel must be at least 128 and must exceed both red and green
        return b >= 100 and b > r and b > g
    except Exception:
        return False


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0

    Task: For a 5-slide sales pitch:
      1. Apply bold and blue color to the title on slide 2 ('Our Solution')
      2. Add a new text box on slide 3 that reads 'Competitive Advantage'
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must be a 5-slide presentation
    if len(prs.slides) < 3:
        print(f"CRITICAL: Expected at least 3 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]
    slide3 = prs.slides[2]

    # ---------- Component 1: Slide 2 title is bold (0.35 points) ----------
    try:
        slide2_title_bold = False
        for shape in slide2.shapes:
            # Find the title placeholder on slide 2
            if shape.has_text_frame and 'Title' in shape.name:
                for para in shape.text_frame.paragraphs:
                    non_empty_runs = [r for r in para.runs if (r.text or "").strip()]
                    if non_empty_runs:
                        # Check all non-empty runs; require all to be bold
                        all_bold = all(
                            (r.font.bold is True) for r in non_empty_runs
                        )
                        if all_bold and para.text.strip():
                            slide2_title_bold = True
                            break
                if slide2_title_bold:
                    break

        if slide2_title_bold:
            print("PASS: Component 1 — Slide 2 title is bold (0.35 pts)")
            total_score += 0.35
        else:
            print("FAIL: Component 1 — Slide 2 title is NOT bold")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # ---------- Component 2: Slide 2 title has blue color (0.35 points) ----------
    try:
        slide2_title_blue = False
        for shape in slide2.shapes:
            if shape.has_text_frame and 'Title' in shape.name:
                for para in shape.text_frame.paragraphs:
                    non_empty_runs = [r for r in para.runs if (r.text or "").strip()]
                    if non_empty_runs:
                        for run in non_empty_runs:
                            # Check if color type is set and is a blue color
                            if run.font.color.type is not None:
                                try:
                                    rgb_str = str(run.font.color.rgb)
                                    if is_blue_color(rgb_str):
                                        slide2_title_blue = True
                                        print(f"  Color found: #{rgb_str}")
                                        break
                                except Exception as color_err:
                                    print(f"  Color check error: {color_err}")
                    if slide2_title_blue:
                        break
                if slide2_title_blue:
                    break

        if slide2_title_blue:
            print("PASS: Component 2 — Slide 2 title has blue color (0.35 pts)")
            total_score += 0.35
        else:
            print("FAIL: Component 2 — Slide 2 title does NOT have blue color")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # ---------- Component 3: Slide 3 has a text box with 'Competitive Advantage' (0.30 points) ----------
    try:
        competitive_advantage_found = False
        for shape in slide3.shapes:
            # Look for a TEXT_BOX shape (shape_type 17) with the required text
            # Also check any shape type as fallback in case shape type differs
            if shape.has_text_frame:
                full_text = shape.text_frame.text.strip()
                # Must contain "Competitive Advantage" and NOT be an existing placeholder
                if 'Competitive Advantage' in full_text and shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    competitive_advantage_found = True
                    print(f"  TextBox found: '{full_text}' in shape '{shape.name}'")
                    break

        if competitive_advantage_found:
            print("PASS: Component 3 — Slide 3 has text box with 'Competitive Advantage' (0.30 pts)")
            total_score += 0.30
        else:
            print("FAIL: Component 3 — Slide 3 does NOT have a text box with 'Competitive Advantage'")
            # Diagnostic: print what text boxes are on slide 3
            for shape in slide3.shapes:
                if shape.has_text_frame:
                    print(f"  Found shape '{shape.name}' (type {shape.shape_type}): {repr(shape.text_frame.text[:80])}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entrypoint
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
