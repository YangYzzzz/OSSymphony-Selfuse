"""
Initial Setup: Product catalog with compressed/blurry product photos on slides 4-8
Task ID: impress_fix_007
Domain: libreoffice_impress

Creates a presentation with 8 slides. Slides 1-3 are title/overview/intro.
Slides 4-8 each contain a low-quality compressed product image.
High-resolution originals are placed in ~/Desktop/product_images/.
"""

import os
import io
import shlex
import subprocess
import time
from PIL import Image, ImageDraw, ImageFont

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_007'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
IMG_DIR = f'{WORKDIR}/Desktop/product_images'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_product_image(filename, product_name, color, width=1600, height=1200, quality='high'):
    """
    Create a product-like image. High quality = large, sharp.
    Low quality = small, heavily compressed JPEG with artifacts.
    """
    img = Image.new('RGB', (width, height), color)
    draw = ImageDraw.Draw(img)

    # Draw some product-like shapes
    # A centered rectangle (product box)
    box_x1 = width // 4
    box_y1 = height // 4
    box_x2 = 3 * width // 4
    box_y2 = 3 * height // 4
    draw.rectangle([box_x1, box_y1, box_x2, box_y2], fill='white', outline='black', width=3)

    # Product name text
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 48)
    except (OSError, IOError):
        font = ImageFont.load_default()

    bbox = draw.textbbox((0, 0), product_name, font=font)
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]
    text_x = (width - tw) // 2
    text_y = (height - th) // 2
    draw.text((text_x, text_y), product_name, fill='black', font=font)

    # Add some decorative elements
    draw.ellipse([width//8, height//8, width//4, height//4], fill=(200, 200, 200), outline='gray')
    draw.ellipse([3*width//4, 3*height//4, 7*width//8, 7*height//8], fill=(200, 200, 200), outline='gray')

    # Diagonal stripe pattern for visual interest
    for i in range(0, width + height, 80):
        draw.line([(i, 0), (i - height, height)], fill=(230, 230, 230), width=1)

    if quality == 'high':
        img.save(filename, 'JPEG', quality=95)
    else:
        # Create heavily compressed, low-quality version
        # First downscale, then upscale to create blur + artifacts
        small = img.resize((width // 6, height // 6), Image.NEAREST)
        blurry = small.resize((width, height), Image.NEAREST)
        blurry.save(filename, 'JPEG', quality=5)

    return filename


def create_initial():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # Product data
    products = [
        {"name": "Aether Pro Wireless Headphones", "color": (45, 85, 140), "price": "$349.99"},
        {"name": "Nova Smart Watch Elite", "color": (75, 130, 65), "price": "$599.00"},
        {"name": "Pulse Fitness Tracker Band", "color": (160, 55, 55), "price": "$129.95"},
        {"name": "Zenith Portable Speaker", "color": (120, 80, 150), "price": "$249.00"},
        {"name": "Lumen LED Desk Lamp", "color": (180, 130, 40), "price": "$89.50"},
    ]

    # Create high-res originals in ~/Desktop/product_images/
    os.makedirs(IMG_DIR, exist_ok=True)
    hq_paths = []
    lq_paths = []

    for i, prod in enumerate(products, 1):
        hq_path = f'{IMG_DIR}/product_{i}.jpg'
        create_product_image(hq_path, prod["name"], prod["color"], quality='high')
        hq_paths.append(hq_path)
        print(f'Created high-res image: {hq_path}')

        # Create low-quality version in a temp location
        lq_path = f'/tmp/product_{i}_low.jpg'
        create_product_image(lq_path, prod["name"], prod["color"], quality='low')
        lq_paths.append(lq_path)

    # Build presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "TechVista Product Catalog"
    slide1.placeholders[1].text = "Spring 2025 Collection"
    # Background
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    # Style title
    for run in slide1.shapes.title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.size = Pt(44)
    for run in slide1.placeholders[1].text_frame.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)

    # --- Slide 2: Table of Contents ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Product Overview"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2D, 0x55, 0x8C)

    # List of products
    for i, prod in enumerate(products, 1):
        txItem = slide2.shapes.add_textbox(Inches(1.5), Inches(1.5 + i * 0.8), Inches(9), Inches(0.6))
        tf2 = txItem.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"{i}. {prod['name']} — {prod['price']}"
        r2 = p2.runs[0]
        r2.font.size = Pt(20)
        r2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 3: Introduction ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox3 = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "About Our Products"
    p3.runs[0].font.size = Pt(36)
    p3.runs[0].font.bold = True
    p3.runs[0].font.color.rgb = RGBColor(0x2D, 0x55, 0x8C)

    body3 = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf3b = body3.text_frame
    tf3b.word_wrap = True
    p3b = tf3b.paragraphs[0]
    p3b.text = (
        "At TechVista, we design cutting-edge consumer electronics that blend "
        "innovation with everyday usability. Each product in our Spring 2025 "
        "collection has been rigorously tested for durability, performance, and "
        "aesthetic appeal. Browse our catalog to discover the perfect addition "
        "to your tech ecosystem."
    )
    p3b.runs[0].font.size = Pt(18)
    p3b.runs[0].font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # --- Slides 4-8: Product slides with LOW QUALITY images ---
    img_left = Inches(1.5)
    img_top = Inches(1.8)
    img_width = Inches(5)
    img_height = Inches(4.5)

    for i, prod in enumerate(products):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank

        # Product title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        ttf = title_box.text_frame
        tp = ttf.paragraphs[0]
        tp.text = prod["name"]
        tp.alignment = PP_ALIGN.CENTER
        tr = tp.runs[0]
        tr.font.size = Pt(32)
        tr.font.bold = True
        tr.font.color.rgb = RGBColor(0x2D, 0x55, 0x8C)

        # LOW QUALITY product image
        pic = slide.shapes.add_picture(
            lq_paths[i],
            img_left, img_top,
            img_width, img_height
        )

        # Price tag
        price_box = slide.shapes.add_textbox(Inches(7.5), Inches(2.5), Inches(4), Inches(1))
        ptf = price_box.text_frame
        pp = ptf.paragraphs[0]
        pp.text = prod["price"]
        pp.alignment = PP_ALIGN.CENTER
        pr = pp.runs[0]
        pr.font.size = Pt(40)
        pr.font.bold = True
        pr.font.color.rgb = RGBColor(0x00, 0x80, 0x00)

        # Description
        desc_box = slide.shapes.add_textbox(Inches(7.5), Inches(3.8), Inches(4.5), Inches(3))
        dtf = desc_box.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        descriptions = [
            "Premium wireless headphones with active noise cancellation, 40-hour battery life, and ultra-comfortable memory foam ear cushions.",
            "Advanced smartwatch with health monitoring, GPS tracking, water resistance up to 50m, and a stunning AMOLED display.",
            "Lightweight fitness tracker with heart rate monitoring, sleep analysis, step counting, and 14-day battery life.",
            "Compact Bluetooth speaker delivering 360-degree sound, IPX7 waterproof rating, and 20-hour playtime.",
            "Adjustable LED desk lamp with 5 color temperature modes, USB charging port, and touch-sensitive controls."
        ]
        dp.text = descriptions[i]
        dp.runs[0].font.size = Pt(16)
        dp.runs[0].font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Clean up temp low-quality images
    for lq in lq_paths:
        if os.path.exists(lq):
            os.remove(lq)

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
