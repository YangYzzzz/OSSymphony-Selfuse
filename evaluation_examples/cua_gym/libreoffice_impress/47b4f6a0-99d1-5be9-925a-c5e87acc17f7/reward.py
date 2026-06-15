"""
Reward Script: Re-insert high-resolution product images on slides 4-8
Task ID: impress_fix_007
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): Images on slides 4-8 are high-resolution (>100KB each, 0.1 per slide)
  Component 2 (0.3): Images match original product files byte-for-byte (0.06 per slide)
  Component 3 (0.2): Image positions/dimensions preserved AND high-res (0.04 per slide)
"""

import os
import hashlib

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_007'

# Expected original image hashes (from Desktop/product_images/product_N.jpg)
EXPECTED_HASHES = {
    3: '019cdfe510e9706f550b72407b9cd9cd',  # slide index 3 = slide 4 -> product_1.jpg
    4: '67429617ccede05f5b0574ece7cb2f20',  # slide index 4 = slide 5 -> product_2.jpg
    5: 'f3aa323bdb5cda3ca3ada45134ee8867',  # slide index 5 = slide 6 -> product_3.jpg
    6: '8c092f7f9f11f246dee36eb0c7a6d244',  # slide index 6 = slide 7 -> product_4.jpg
    7: '27572401bab9bb5e38ea9b12abf3d116',  # slide index 7 = slide 8 -> product_5.jpg
}

# Expected positions and sizes (EMU) for images on slides 4-8
EXPECTED_LEFT = 1371600
EXPECTED_TOP = 1645920
EXPECTED_WIDTH = 4572000
EXPECTED_HEIGHT = 4114800
POSITION_TOLERANCE = 0.01  # 1% relative tolerance


def is_approx_equal(val1, val2, tol=POSITION_TOLERANCE):
    """Check approximate equality with relative tolerance."""
    if val1 == val2:
        return True
    if val1 == 0 or val2 == 0:
        return abs(val1 - val2) < 10000  # small absolute tolerance for zero
    return abs(val1 - val2) / max(abs(val1), abs(val2)) <= tol


def get_picture_shape(slide):
    """Find the first PICTURE shape on a slide."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return shape
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 8 slides
    if len(prs.slides) < 8:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 8")
        print("REWARD: 0.0")
        return 0.0

    # Check each slide 4-8 (indices 3-7)
    for slide_idx in range(3, 8):
        slide_num = slide_idx + 1
        slide = prs.slides[slide_idx]
        pic = get_picture_shape(slide)

        if pic is None:
            print(f"FAIL: Slide {slide_num} has no picture shape")
            continue

        blob = pic.image.blob
        blob_size = len(blob)
        blob_md5 = hashlib.md5(blob).hexdigest()
        expected_md5 = EXPECTED_HASHES[slide_idx]

        # Component 1: Image is high-resolution (>100KB) — 0.1 points per slide
        try:
            if blob_size > 100000:
                print(f"PASS: Comp1 — Slide {slide_num} image is high-res ({blob_size} bytes) (0.1 pts)")
                total_score += 0.1
            else:
                print(f"FAIL: Comp1 — Slide {slide_num} image is low-res ({blob_size} bytes, need >100KB)")
        except Exception as e:
            print(f"ERROR: Comp1 — Slide {slide_num}: {e}")

        # Component 2: Image matches original product file — 0.06 points per slide
        try:
            if blob_md5 == expected_md5:
                print(f"PASS: Comp2 — Slide {slide_num} image matches original (md5={blob_md5[:12]}...) (0.06 pts)")
                total_score += 0.06
            else:
                print(f"FAIL: Comp2 — Slide {slide_num} md5={blob_md5[:12]}... expected={expected_md5[:12]}...")
        except Exception as e:
            print(f"ERROR: Comp2 — Slide {slide_num}: {e}")

        # Component 3: Position/size correct AND image is high-res — 0.04 points per slide
        try:
            pos_ok = (
                is_approx_equal(pic.left, EXPECTED_LEFT) and
                is_approx_equal(pic.top, EXPECTED_TOP) and
                is_approx_equal(pic.width, EXPECTED_WIDTH) and
                is_approx_equal(pic.height, EXPECTED_HEIGHT)
            )
            if pos_ok and blob_size > 100000:
                print(f"PASS: Comp3 — Slide {slide_num} position/size correct AND high-res (0.04 pts)")
                total_score += 0.04
            elif pos_ok:
                print(f"FAIL: Comp3 — Slide {slide_num} position correct but image still low-res")
            else:
                print(f"FAIL: Comp3 — Slide {slide_num} position/size mismatch: "
                      f"L={pic.left} T={pic.top} W={pic.width} H={pic.height}")
        except Exception as e:
            print(f"ERROR: Comp3 — Slide {slide_num}: {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
