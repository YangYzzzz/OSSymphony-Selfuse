"""
Initial Setup: Sales transformation presentation with 8 slides
Task ID: impress_sales_066
Domain: libreoffice_impress

Creates Transform_Pitch.pptx with 8 slides. Slide 4 has only the title
'The Transformation' (no before/after content yet). Also creates
before.png and after.png on the Desktop.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# --- Generate placeholder images using PIL ---
from PIL import Image, ImageDraw, ImageFont

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
TASK_ID = 'impress_sales_066'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_before_image():
    """Create a cluttered spreadsheet-style image."""
    img = Image.new('RGB', (640, 480), '#F0F0F0')
    draw = ImageDraw.Draw(img)
    # Draw grid lines to look like a messy spreadsheet
    for y in range(0, 480, 20):
        draw.line([(0, y), (640, y)], fill='#CCCCCC', width=1)
    for x in range(0, 640, 80):
        draw.line([(x, 0), (x, 480)], fill='#CCCCCC', width=1)
    # Add some "data" rectangles in various colors to look cluttered
    colors = ['#FFD700', '#FF6347', '#4682B4', '#90EE90', '#DDA0DD', '#FFA07A']
    import random
    random.seed(42)
    for i in range(30):
        x1 = random.randint(5, 560)
        y1 = random.randint(5, 400)
        x2 = x1 + random.randint(40, 80)
        y2 = y1 + random.randint(15, 25)
        draw.rectangle([x1, y1, x2, y2], fill=colors[i % len(colors)], outline='#888888')
    # Title bar
    draw.rectangle([0, 0, 640, 30], fill='#2E4057')
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 14)
    except:
        font = ImageFont.load_default()
    draw.text((10, 8), "Quarterly_Data_ALL_v17_FINAL_revised.xlsx", fill='white', font=font)
    os.makedirs(DESKTOP, exist_ok=True)
    img.save(f'{DESKTOP}/before.png')
    print(f'Created {DESKTOP}/before.png')


def create_after_image():
    """Create a clean dashboard-style image."""
    img = Image.new('RGB', (640, 480), '#FFFFFF')
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 16)
        font_sm = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 12)
    except:
        font = ImageFont.load_default()
        font_sm = font

    # Header bar
    draw.rectangle([0, 0, 640, 50], fill='#1B3A5C')
    draw.text((20, 15), "Sales Performance Dashboard", fill='white', font=font)

    # KPI cards
    kpi_data = [
        ("Revenue", "$2.4M", '#27AE60'),
        ("Growth", "+18%", '#2980B9'),
        ("Customers", "1,247", '#8E44AD'),
    ]
    for i, (label, value, color) in enumerate(kpi_data):
        x = 30 + i * 200
        draw.rounded_rectangle([x, 70, x + 170, 150], radius=8, fill=color)
        draw.text((x + 15, 85), label, fill='white', font=font_sm)
        draw.text((x + 15, 110), value, fill='white', font=font)

    # Chart area placeholder
    draw.rectangle([30, 170, 610, 390], outline='#DDDDDD', width=2)
    # Simple bar chart
    bars = [180, 220, 195, 260, 240, 290, 310, 280, 330, 350, 320, 370]
    bar_w = 40
    for i, h in enumerate(bars):
        x = 50 + i * 46
        y_top = 380 - (h - 150)
        draw.rectangle([x, y_top, x + bar_w, 380], fill='#3498DB')
    # X-axis labels
    months = ['J', 'F', 'M', 'A', 'M', 'J', 'J', 'A', 'S', 'O', 'N', 'D']
    for i, m in enumerate(months):
        draw.text((58 + i * 46, 385), m, fill='#666666', font=font_sm)

    # Footer
    draw.text((30, 420), "Last updated: March 2025  |  Automated refresh: Daily", fill='#999999', font=font_sm)

    os.makedirs(DESKTOP, exist_ok=True)
    img.save(f'{DESKTOP}/after.png')
    print(f'Created {DESKTOP}/after.png')


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a styled textbox."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Sales Transformation Strategy"
    slide1.placeholders[1].text = "Q2 2025 Executive Briefing\nPrepared by the Revenue Operations Team"

    # --- Slide 2: Market Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Market Overview"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Global SaaS market projected to reach $307B by 2026"
    p2 = tf2.add_paragraph()
    p2.text = "Enterprise segment growing at 14.2% CAGR"
    p2.level = 1
    p3 = tf2.add_paragraph()
    p3.text = "Mid-market adoption accelerating with AI-driven tools"
    p3.level = 1
    p4 = tf2.add_paragraph()
    p4.text = "Key competitors: Salesforce, HubSpot, Pipedrive"
    p4.level = 1
    p5 = tf2.add_paragraph()
    p5.text = "Our market share: 4.7% (up from 3.2% in 2023)"
    p5.level = 1

    # --- Slide 3: Current Challenges ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Current Challenges"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Sales cycle averaging 47 days (industry benchmark: 32 days)"
    p = tf3.add_paragraph()
    p.text = "CRM data quality issues affecting 38% of records"
    p.level = 1
    p = tf3.add_paragraph()
    p.text = "Manual reporting consuming 12+ hours per week per manager"
    p.level = 1
    p = tf3.add_paragraph()
    p.text = "Pipeline visibility limited to weekly snapshots"
    p.level = 1
    p = tf3.add_paragraph()
    p.text = "Win rate declined from 28% to 23% over past two quarters"
    p.level = 1

    # --- Slide 4: The Transformation (title only, no content yet) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    add_textbox(slide4, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                "The Transformation", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C), alignment=PP_ALIGN.CENTER)

    # --- Slide 5: Solution Pipeline ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Solution Pipeline"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Phase 1: CRM data cleanup and migration (April - May)"
    p = tf5.add_paragraph()
    p.text = "Phase 2: Automated reporting dashboard deployment (June)"
    p.level = 0
    p = tf5.add_paragraph()
    p.text = "Phase 3: AI-powered lead scoring integration (July - August)"
    p.level = 0
    p = tf5.add_paragraph()
    p.text = "Phase 4: Real-time pipeline analytics (September)"
    p.level = 0

    # --- Slide 6: Implementation Timeline ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Implementation Timeline"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Week 1-2: Stakeholder interviews and requirements gathering"
    p = tf6.add_paragraph()
    p.text = "Week 3-4: Data audit and cleansing protocols"
    p.level = 1
    p = tf6.add_paragraph()
    p.text = "Week 5-8: Platform configuration and integration testing"
    p.level = 1
    p = tf6.add_paragraph()
    p.text = "Week 9-10: User training and pilot launch"
    p.level = 1
    p = tf6.add_paragraph()
    p.text = "Week 11-12: Full rollout and performance monitoring"
    p.level = 1

    # --- Slide 7: Team & Resources ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Team & Resources"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "Project Lead: Sarah Chen, VP Revenue Operations"
    p = tf7.add_paragraph()
    p.text = "Technical Lead: Marcus Johnson, Director of Sales Engineering"
    p.level = 0
    p = tf7.add_paragraph()
    p.text = "Data Analyst: Priya Patel, Senior Business Analyst"
    p.level = 0
    p = tf7.add_paragraph()
    p.text = "Budget: $420,000 (approved by CFO)"
    p.level = 0
    p = tf7.add_paragraph()
    p.text = "Expected ROI: 340% over 18 months"
    p.level = 0

    # --- Slide 8: Q&A ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    add_textbox(slide8, Inches(2), Inches(2.5), Inches(9), Inches(2),
                "Questions & Discussion", font_size=40, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C), alignment=PP_ALIGN.CENTER)
    add_textbox(slide8, Inches(3), Inches(4.5), Inches(7), Inches(1),
                "Contact: transformation@company.com", font_size=18,
                bold=False, color=RGBColor(0x66, 0x66, 0x66),
                alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Presentation created: {OUTPUT} ({len(prs.slides)} slides)')


if __name__ == '__main__':
    create_before_image()
    create_after_image()
    create_presentation()

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')
