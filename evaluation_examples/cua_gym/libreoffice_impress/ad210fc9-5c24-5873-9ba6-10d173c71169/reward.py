"""
Reward Script: Before/After Transformation Slide on Slide 4
Task ID: impress_sales_066
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.30): "Before" section — label, image, 3 red bullets
  - Component 2 (0.30): "After" section — label, image, 3 green bullets
  - Component 3 (0.15): Vertical divider near horizontal center
  - Component 4 (0.15): Wipe transition on slide 4
  - Component 5 (0.10): Correct images (before.png = 7627 bytes, after.png = 15496 bytes)
"""

import os
import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_066'


def persist_app_state(domain):
    """Save any unsaved LibreOffice changes."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print(f"PERSIST: ctrl+s sent for {domain}")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def get_all_text_shapes(slide):
    """Get all shapes with text frames, including nested in groups."""
    def extract(shape):
        results = []
        if hasattr(shape, "text") and hasattr(shape, "text_frame"):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.extend(extract(sub))
        return results
    out = []
    for shape in slide.shapes:
        out.extend(extract(shape))
    return out


def check_color_match(run, expected_hex):
    """Check if a run's font color matches the expected hex (e.g. 'CC0000')."""
    try:
        if run.font.color.type is not None:
            actual = str(run.font.color.rgb).upper()
            return actual == expected_hex.upper()
    except Exception:
        pass
    return False


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)
    slide_width = prs.slide_width

    # Collect all shapes on slide 4
    all_shapes = list(slide.shapes)
    text_shapes = get_all_text_shapes(slide)
    picture_shapes = [s for s in all_shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]

    # ------------------------------------------------------------------
    # Component 1: "Before" section (0.30 points)
    #   - A text shape with "Before" label in red (#CC0000)
    #   - 3 red (#CC0000) bullet text items
    #   - An image on the left half of the slide
    # ------------------------------------------------------------------
    try:
        before_label_found = False
        red_bullets_count = 0

        for shape in text_shapes:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                # Check for "Before" label
                for run in para.runs:
                    if run.text.strip().lower() == 'before' and check_color_match(run, 'CC0000'):
                        before_label_found = True
                        break
                # Check for red bullet points (non-label text in CC0000)
                if text.lower() != 'before' and len(text) > 10:
                    for run in para.runs:
                        if check_color_match(run, 'CC0000'):
                            red_bullets_count += 1
                            break

        # Check for an image in the left half
        left_image_found = False
        for pic in picture_shapes:
            pic_center_x = pic.left + pic.width / 2
            if pic_center_x < slide_width / 2:
                left_image_found = True
                break

        before_score = 0.0
        if before_label_found:
            before_score += 0.10
            print(f"PASS: 'Before' label found with red color CC0000")
        else:
            print(f"FAIL: 'Before' label not found or not red CC0000")

        if red_bullets_count >= 3:
            before_score += 0.10
            print(f"PASS: Found {red_bullets_count} red bullet points (need >= 3)")
        else:
            print(f"FAIL: Found {red_bullets_count} red bullet points (need >= 3)")

        if left_image_found:
            before_score += 0.10
            print(f"PASS: Image found in left half of slide")
        else:
            print(f"FAIL: No image found in left half of slide")

        total_score += before_score
    except Exception as e:
        print(f"ERROR: Component 1 (Before section) — {e}")

    # ------------------------------------------------------------------
    # Component 2: "After" section (0.30 points)
    #   - A text shape with "After" label in green (#00AA00)
    #   - 3 green (#00AA00) bullet text items
    #   - An image on the right half of the slide
    # ------------------------------------------------------------------
    try:
        after_label_found = False
        green_bullets_count = 0

        for shape in text_shapes:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                for run in para.runs:
                    if run.text.strip().lower() == 'after' and check_color_match(run, '00AA00'):
                        after_label_found = True
                        break
                if text.lower() != 'after' and len(text) > 10:
                    for run in para.runs:
                        if check_color_match(run, '00AA00'):
                            green_bullets_count += 1
                            break

        # Check for an image in the right half
        right_image_found = False
        for pic in picture_shapes:
            pic_center_x = pic.left + pic.width / 2
            if pic_center_x > slide_width / 2:
                right_image_found = True
                break

        after_score = 0.0
        if after_label_found:
            after_score += 0.10
            print(f"PASS: 'After' label found with green color 00AA00")
        else:
            print(f"FAIL: 'After' label not found or not green 00AA00")

        if green_bullets_count >= 3:
            after_score += 0.10
            print(f"PASS: Found {green_bullets_count} green bullet points (need >= 3)")
        else:
            print(f"FAIL: Found {green_bullets_count} green bullet points (need >= 3)")

        if right_image_found:
            after_score += 0.10
            print(f"PASS: Image found in right half of slide")
        else:
            print(f"FAIL: No image found in right half of slide")

        total_score += after_score
    except Exception as e:
        print(f"ERROR: Component 2 (After section) — {e}")

    # ------------------------------------------------------------------
    # Component 3: Vertical divider near center (0.15 points)
    #   - A narrow shape (width < 1 inch = 914400 EMU) that is tall
    #   - Positioned near the horizontal center of the slide
    # ------------------------------------------------------------------
    try:
        divider_found = False
        for shape in all_shapes:
            # A divider is either an AUTO_SHAPE or a line that is narrow and tall
            if shape.width < 914400 and shape.height > 2000000:
                # Check if roughly centered (within 15% of slide center)
                shape_center_x = shape.left + shape.width / 2
                slide_center_x = slide_width / 2
                offset_ratio = abs(shape_center_x - slide_center_x) / slide_width
                if offset_ratio < 0.15:
                    divider_found = True
                    print(f"PASS: Vertical divider found — name={shape.name}, center_x={shape_center_x}, slide_center={slide_center_x}, offset={offset_ratio:.3f}")
                    break

        if divider_found:
            total_score += 0.15
        else:
            print(f"FAIL: No vertical divider found near center of slide 4")
    except Exception as e:
        print(f"ERROR: Component 3 (Divider) — {e}")

    # ------------------------------------------------------------------
    # Component 4: Wipe transition on slide 4 (0.15 points)
    #   - Slide 4 must have a <p:transition> with a <p:wipe> child element
    # ------------------------------------------------------------------
    try:
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        wipe_found = False
        with zipfile.ZipFile(file_path, 'r') as zf:
            try:
                with zf.open('ppt/slides/slide4.xml') as f:
                    root = ET.parse(f).getroot()
                    trans = root.find('.//p:transition', ns)
                    if trans is not None:
                        wipe_elem = trans.find('.//p:wipe', ns)
                        if wipe_elem is not None:
                            wipe_found = True
                            print(f"PASS: Wipe transition found on slide 4")
                        else:
                            # Check any child for 'wipe' in tag name
                            for child in trans:
                                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                                if 'wipe' in tag.lower():
                                    wipe_found = True
                                    print(f"PASS: Wipe transition found on slide 4 (tag: {tag})")
                                    break
                            if not wipe_found:
                                children = [c.tag.split('}')[-1] if '}' in c.tag else c.tag for c in trans]
                                print(f"FAIL: Transition exists but no Wipe — found: {children}")
                    else:
                        print(f"FAIL: No transition element on slide 4")
            except KeyError:
                print(f"FAIL: Could not open slide4.xml")

        if wipe_found:
            total_score += 0.15
    except Exception as e:
        print(f"ERROR: Component 4 (Wipe transition) — {e}")

    # ------------------------------------------------------------------
    # Component 5: Correct images used (0.10 points)
    #   - before.png (7627 bytes) on the left, after.png (15496 bytes) on the right
    #   - Verify by blob size: left image blob ~7627 bytes, right image blob ~15496 bytes
    # ------------------------------------------------------------------
    try:
        images_correct = False
        left_blob_len = None
        right_blob_len = None

        for pic in picture_shapes:
            pic_center_x = pic.left + pic.width / 2
            blob_len = len(pic.image.blob)
            if pic_center_x < slide_width / 2:
                left_blob_len = blob_len
            else:
                right_blob_len = blob_len

        # before.png = 7627 bytes blob, after.png = 15496 bytes blob
        if left_blob_len == 7627 and right_blob_len == 15496:
            images_correct = True
            print(f"PASS: Correct images — left=before.png ({left_blob_len}B), right=after.png ({right_blob_len}B)")
        else:
            # Allow some tolerance — at least both images present on correct sides
            if left_blob_len is not None and right_blob_len is not None:
                print(f"FAIL: Images present but sizes don't match — left={left_blob_len}B (expected 7627), right={right_blob_len}B (expected 15496)")
            else:
                print(f"FAIL: Missing images — left={left_blob_len}, right={right_blob_len}")

        if images_correct:
            total_score += 0.10
    except Exception as e:
        print(f"ERROR: Component 5 (Images) — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
