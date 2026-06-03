"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 209 feels cramped—the body copy is all squished together. In LibreOffice Impress, how do I adjust the main text box so the line spacing is set to Exactly 20 pt and the paragraphs are fully Justified?
Generated: 2025-09-10 20:34:48
Status: success
Model: azure-o3
Total Steps: 12
"""

import os
from pptx import Presentation
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import PP_PLACEHOLDER as P


def _points_from_line_spacing(line_spacing):
    """Convert python-pptx Paragraph.line_spacing to points (float) when possible."""
    if line_spacing is None:
        return None
    # python-pptx returns EMU (1 pt = 12 700 EMU) for ‘Exactly’ spacing
    if line_spacing > 1000:  # assume EMU value
        return line_spacing / 12_700.0
    return float(line_spacing)


def _is_main_body_placeholder(shape):
    """Return True if the shape is a body/content placeholder (not title/subtitle)."""
    if not shape.is_placeholder:
        return False
    ph_type = shape.placeholder_format.type
    return ph_type not in (P.TITLE, P.CENTER_TITLE, P.SUBTITLE)


def verify_task(file_path: str, target_slide_number: int = 209) -> float:
    """Verify that on the target slide the body text paragraphs
    1) are fully justified and
    2) have line spacing set to Exactly 20 pt.
    Returns a progressive score between 0.0 and 1.0.
    """

    print("Checking task completion for line spacing & justification on slide 209…")

    # -------- File existence & loading (no score for mere existence) -------- #
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0
    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Unable to open presentation: {e}")
        return 0.0

    # -------- Locate target slide -------- #
    slide_index = target_slide_number - 1  # zero-based index
    if slide_index >= len(prs.slides):
        print(f"! Slide {target_slide_number} not found – using last slide as fallback")
        slide_index = len(prs.slides) - 1
    slide = prs.slides[slide_index]

    # -------- Analyse body paragraphs -------- #
    total_paragraphs = 0
    justified_ok = 0
    spacing_ok = 0

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if not _is_main_body_placeholder(shape):  # focus on main text box
            continue

        for para in shape.text_frame.paragraphs:
            text = ''.join(run.text for run in para.runs).strip()
            if not text:
                continue  # ignore empty paragraphs

            total_paragraphs += 1

            # Check justification
            if para.alignment == PP_PARAGRAPH_ALIGNMENT.JUSTIFY:
                justified_ok += 1

            # Check line spacing ≈ 20 pt (tolerance 0.5 pt)
            pts = _points_from_line_spacing(para.line_spacing)
            if pts is not None and abs(pts - 20) <= 0.5:
                spacing_ok += 1

    if total_paragraphs == 0:
        print("✗ No body text paragraphs found on target slide")
        return 0.0

    # -------- Scoring (progressive) -------- #
    justify_ratio = justified_ok / total_paragraphs
    spacing_ratio = spacing_ok / total_paragraphs

    print(f"Total body paragraphs: {total_paragraphs}")
    print(f"  • Justified: {justified_ok} ({justify_ratio:.0%})")
    print(f"  • Exactly 20 pt line spacing: {spacing_ok} ({spacing_ratio:.0%})")

    # 0.5 for each requirement
    score = 0.5 * justify_ratio + 0.5 * spacing_ratio
    final_score = round(min(score, 1.0), 3)

    print(f"REWARD: {final_score}")
    return final_score


# ---------------- Execute verification when run as script ---------------- #
if __name__ == "__main__":
    FILE_PATH = "/home/user/slide_209_feels_crampedthe_body_copy_is_all_squished_together_in_libreoffice_impress_how_do_i_adjust_golden.pptx"
    verify_task(FILE_PATH)
