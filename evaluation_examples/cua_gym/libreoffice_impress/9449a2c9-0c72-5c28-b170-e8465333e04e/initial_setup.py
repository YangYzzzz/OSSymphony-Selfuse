"""
Initial Setup: Medical conference presentation with clinical notes
Task ID: osworld_multi_apps_impress_notes_export_008
Domain: libreoffice_impress

Creates Cardiology_Symposium.pptx on the Desktop with 11 slides
covering treatment protocols and clinical presenter notes.
cardiology_notes.docx must NOT exist initially.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_impress_notes_export_008'
OUTPUT = f'{WORKDIR}/Cardiology_Symposium.pptx'
# Make sure cardiology_notes.docx does NOT exist
NOTES_FILE = f'{WORKDIR}/cardiology_notes.docx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    # Remove notes file if accidentally exists
    if os.path.exists(NOTES_FILE):
        os.remove(NOTES_FILE)
        print(f'Removed pre-existing notes file: {NOTES_FILE}')

    os.makedirs(WORKDIR, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide data: (title, subtitle_or_body, notes_text)
    slides_data = [
        (
            "2025 International Cardiology Symposium",
            "Advances in Heart Failure Management\nDr. Elena Vasquez, MD, PhD | March 2025",
            "Welcome delegates to the 2025 International Cardiology Symposium.\nToday we present the latest evidence-based protocols for heart failure management.\nThis talk covers both HFrEF and HFpEF subgroups.\nPlease hold questions until the end of each section."
        ),
        (
            "Heart Failure: Epidemiology & Burden",
            "• 64 million patients worldwide affected\n• 5-year mortality rate: 50% in advanced HF\n• $108 billion annual healthcare cost (US)\n• 1.9 million hospitalizations per year (US)",
            "Heart failure remains a leading cause of morbidity and mortality globally.\nNote the disproportionate burden on elderly patients over 75.\nReadmission rates remain above 25% within 30 days despite improvements.\nEmphasize the social determinants of health that influence outcomes.\nHospitalization costs account for roughly 60% of total HF expenditure."
        ),
        (
            "Classification & Staging",
            "NYHA Functional Classification:\n• Class I: No symptoms\n• Class II: Slight limitation\n• Class III: Marked limitation\n• Class IV: Symptoms at rest\n\nACC/AHA Staging A–D",
            "The NYHA classification guides symptom assessment and treatment intensity.\nACC/AHA staging emphasizes progressive structural changes.\nStage A includes risk factors but no structural disease—focus on prevention.\nStage D represents refractory HF requiring advanced therapies.\nCombining both systems gives a fuller clinical picture."
        ),
        (
            "HFrEF: First-Line Pharmacotherapy",
            "GDMT Pillars:\n1. ACEi / ARB / ARNI\n   - Sacubitril/Valsartan (PARADIGM-HF)\n2. Beta-blockers\n   - Carvedilol, Metoprolol Succinate, Bisoprolol\n3. MRA\n   - Spironolactone, Eplerenone\n4. SGLT2i\n   - Dapagliflozin, Empagliflozin",
            "GDMT stands for Guideline-Directed Medical Therapy.\nSacubitril/valsartan demonstrated 20% reduction in CV death vs enalapril in PARADIGM-HF.\nBeta-blocker selection is critical: only carvedilol, metoprolol succinate, and bisoprolol are proven.\nSpironolactone 25mg daily is first-line MRA if GFR > 30 and K+ < 5.0 mEq/L.\nSGLT2 inhibitors are now Class I recommendation regardless of diabetes status.\nTarget doses over 3-6 months; do not rush up-titration in fragile patients."
        ),
        (
            "ARNI Therapy: Sacubitril/Valsartan",
            "Key Trial: PARADIGM-HF\n• 8,442 patients, LVEF ≤ 40%\n• Primary endpoint: CV death or HF hospitalization\n• Relative risk reduction: 20% (NNT = 21)\n• Absolute risk reduction: 4.7%\n\nInitiation: Wash-out 36 hrs from ACEi\nStarting dose: 24/26 mg BID → target 97/103 mg BID",
            "PARADIGM-HF was a landmark trial published in NEJM 2014.\nTrial was stopped early due to overwhelming benefit.\nKey safety concern: angioedema risk, especially in African-American patients.\nMandatory 36-hour washout from ACEi before initiating sacubitril/valsartan.\nAvoid in patients with history of angioedema.\nMonitor blood pressure, potassium, and renal function at each dose step.\nCost-effectiveness analyses support its use when generic ARNIs become available."
        ),
        (
            "Device Therapy in HFrEF",
            "ICD Indications:\n• LVEF ≤ 35%, NYHA II-III, >1 yr survival\n• Post-MI >40 days\n\nCRT Indications:\n• LVEF ≤ 35%, QRS ≥ 150ms (LBBB)\n• NYHA II-IV on GDMT\n\nLVAD as Bridge or Destination Therapy",
            "Device therapy is additive to GDMT, not a substitute.\nImportant: ICD benefit requires LVEF reassessment after 3 months of optimal GDMT.\nCRT improves symptoms, exercise capacity, and reduces hospitalization in responders.\nApproximately 30% of CRT patients are non-responders—patient selection is key.\nLVAD destination therapy (INTERMACS profile ≥ 2) now has 2-year survival over 70%.\nShared decision-making is essential when discussing device therapy goals."
        ),
        (
            "HFpEF: Emerging Evidence",
            "• Prevalence: >50% of all HF cases\n• Pathophysiology: Diastolic dysfunction, inflammation\n• No proven mortality-reducing therapy until EMPEROR-Preserved\n\nEmpagliflozin: 21% reduction in CV death or HF hospitalization\nSpironolactone: TOPCAT trial—inconclusive but regional data suggest benefit",
            "HFpEF management has historically lacked high-quality evidence.\nEMPEROR-Preserved trial enrolled patients with LVEF > 40%—landmark result.\nEmpagliflozin reduced primary endpoint: HR 0.79 (95% CI 0.69-0.90), p<0.001.\nNote the heterogeneity: patients with LVEF 41-60% (HFmrEF) may benefit most.\nSpironolactone benefit in TOPCAT was driven by Americas subgroup; post-hoc finding.\nBlood pressure control to target <130/80 remains universal recommendation.\nExercise training (supervised cardiac rehab) improves functional capacity in HFpEF."
        ),
        (
            "Acute Decompensated Heart Failure",
            "ADHF Management Priorities:\n1. Diuresis: IV furosemide ≥ 2.5x oral dose\n2. Vasodilators: Nitroglycerin, Nitroprusside\n3. Inotropes: Dobutamine, Milrinone (cardiogenic shock)\n4. Ultrafiltration: Refractory diuretic resistance\n\nDISCHARGE trial: Outpatient IV diuresis noninferior to inpatient in select patients",
            "ADHF treatment focuses on decongestion while avoiding hypoperfusion.\nInitial diuretic dose: IV furosemide at 2.5x the patient's total oral dose (DOSE trial).\nBicarb strategy: alkalosis from over-diuresis worsens outcomes—target Cl-responsive state.\nDobutamine improves hemodynamics but increases arrhythmia risk; use lowest effective dose.\nMilrinone preferred in beta-blocker treated patients; dobutamine may be less effective.\nDISCHARGE trial supports early discharge with IV diuresis follow-up in stable patients.\nClose outpatient follow-up within 7 days is mandatory after ADHF discharge."
        ),
        (
            "Cardiorenal Syndrome",
            "Types I–V:\n• Type 1: Acute HF → Acute Kidney Injury\n• Type 2: Chronic HF → CKD\n• Type 3: Acute Kidney Injury → Acute HF\n• Type 4: CKD → Chronic HF\n• Type 5: Systemic Disorder → Both\n\nKey Biomarkers: BNP, NT-proBNP, Cystatin C, NGAL",
            "Cardiorenal syndrome (CRS) affects up to 40% of hospitalized HF patients.\nType 1 CRS is the most clinically challenging: aggressive diuresis worsens renal function.\nFor CRS Type 1: distinguish true AKI from pre-renal azotemia using FENa.\nNGAL (neutrophil gelatinase-associated lipocalin) is an early marker of tubular injury.\nCystatin C outperforms creatinine in estimating GFR in HF patients.\nAvoiding nephrotoxins (NSAIDs, contrast, aminoglycosides) is paramount.\nRCT data on tolvaptan (V2 receptor antagonist) show decongestion without renal harm."
        ),
        (
            "Palliative Care Integration",
            "Principles:\n• Early integration with HF management\n• Symptom burden assessment: KCCQ score\n• Goals of care discussion: ACP, POLST\n• Deactivation of ICD in end-stage disease\n\nICU-to-Home: Transition pathways for Stage D HF",
            "Palliative care should be introduced at ACC/AHA Stage C, not reserved for Stage D.\nKCCQ (Kansas City Cardiomyopathy Questionnaire) correlates with prognosis and QoL.\nAdvance care planning discussions should include preferences around hospitalization and resuscitation.\nICD deactivation is ethically appropriate in patients with terminal prognosis who consent.\nFamily meetings and chaplaincy involvement reduce moral distress in ICU teams.\nHospice enrollment criteria for HF: NYHA IV, LVEF < 20%, recurrent hospitalizations.\nMedical aid in dying is distinct from palliative care and follows jurisdictional laws."
        ),
        (
            "Future Directions & Conclusions",
            "Emerging Therapies:\n• Mavacamten (HCM, potential HFpEF trials)\n• Omecamtiv Mecarbil (HFrEF - GALACTIC-HF results mixed)\n• RNA therapeutics: TTR-targeted siRNA\n• Remote monitoring: CardioMEMS™ implant\n\nKey Takeaways:\n✓ Optimize GDMT before device therapy\n✓ Reassess LVEF at 3 months\n✓ Integrate palliative care early\n✓ Consider HFpEF in all breathless patients",
            "Mavacamten is FDA-approved for obstructive HCM; trials in HFpEF are ongoing.\nGALACTIC-HF did not show mortality benefit for omecamtiv mecarbil despite improved contractility.\nTTR amyloidosis is an underdiagnosed cause of HFpEF in elderly patients; tafamidis is approved.\nCardioMEMS PA pressure sensor reduces HF hospitalization by 28% in CHAMPION trial.\nRemote monitoring programs require dedicated HF nursing support to realize benefit.\nConclusion: The four-pillar GDMT approach should be prioritized in all eligible HFrEF patients.\nThank the audience and invite questions.\nNext symposium: Berlin, September 2025."
        ),
    ]

    for i, (title, body, notes) in enumerate(slides_data):
        if i == 0:
            # Title slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title
            slide.placeholders[1].text = body
        else:
            # Content slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title
            tf = slide.placeholders[1].text_frame
            tf.text = body

        # Set presenter notes
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = notes

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open the file in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
