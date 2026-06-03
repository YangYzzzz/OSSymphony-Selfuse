"""
Reward Script: Extract presenter notes from Cardiology_Symposium.pptx and save as cardiology_notes.docx
Task ID: osworld_multi_apps_impress_notes_export_008
Domain: libreoffice_impress (multi-app: also involves writing a .docx)
Scoring:
  Component 1: cardiology_notes.docx exists on Desktop (precondition gate — 0.0 if missing)
  Component 2: Paragraph count matches expected notes structure (0.3 pts)
  Component 3: First slide notes content preserved correctly (0.2 pts)
  Component 4: All notes text content matches pptx notes in slide order (0.5 pts)
"""

import os

# python-pptx for reading source presentation
try:
    from pptx import Presentation
except ImportError:
    print("CRITICAL: python-pptx not available")
    print("REWARD: 0.0")
    raise

# python-docx for reading output document
try:
    from docx import Document
except ImportError:
    print("CRITICAL: python-docx not available")
    print("REWARD: 0.0")
    raise

WORKDIR = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_impress_notes_export_008'

PPTX_PATH = f'{WORKDIR}/Cardiology_Symposium.pptx'
DOCX_PATH = f'{WORKDIR}/cardiology_notes.docx'


def build_expected_paragraphs(pptx_path):
    """
    Build the expected list of paragraph texts from the pptx notes.
    Each slide's notes lines are added in order, separated by an empty string
    between slides (matching the docx structure observed in golden).
    Returns a list of strings.
    """
    prs = Presentation(pptx_path)
    expected_lines = []
    for slide in prs.slides:
        if slide.has_notes_slide:
            tf = slide.notes_slide.notes_text_frame
            notes_text = tf.text.strip()
            if notes_text:
                lines = notes_text.split('\n')
                expected_lines.extend(lines)
                expected_lines.append('')  # slide separator blank line

    # Remove trailing blank lines
    while expected_lines and expected_lines[-1] == '':
        expected_lines.pop()

    return expected_lines


def verify_task(docx_path, pptx_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition gate: both files must exist
    if not os.path.exists(docx_path):
        print(f"CRITICAL: Output file not found: {docx_path}")
        print("REWARD: 0.0")
        return 0.0

    if not os.path.exists(pptx_path):
        print(f"CRITICAL: Source file not found: {pptx_path}")
        print("REWARD: 0.0")
        return 0.0

    # Load docx
    try:
        doc = Document(docx_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load docx {docx_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Load expected content from pptx
    try:
        expected_paras = build_expected_paragraphs(pptx_path)
    except Exception as e:
        print(f"CRITICAL: Cannot process pptx {pptx_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    doc_paras = [p.text for p in doc.paragraphs]
    expected_count = len(expected_paras)
    doc_count = len(doc_paras)

    print(f"INFO: Expected paragraph count: {expected_count}")
    print(f"INFO: Actual paragraph count:   {doc_count}")

    # Component 1: Paragraph count matches expected (0.3 points)
    # The docx must have the same number of paragraphs as expected from extracting
    # all notes lines plus one blank separator between each slide's block.
    try:
        if doc_count == expected_count:
            print(f"PASS: Component 1 — Paragraph count matches expected ({expected_count}) (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 — Paragraph count mismatch: expected {expected_count}, found {doc_count}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: First slide notes content preserved correctly (0.2 points)
    # Verifies that slide 1 notes (first 4 lines from pptx) appear correctly in docx
    try:
        # Slide 1 has 4 lines of notes
        slide1_expected = expected_paras[:4]  # first 4 lines before first blank separator
        if doc_count >= 4:
            slide1_actual = doc_paras[:4]
            if slide1_expected == slide1_actual:
                print(f"PASS: Component 2 — Slide 1 notes content correct (0.2 pts)")
                print(f"  First line: {repr(slide1_actual[0][:80])}")
                total_score += 0.2
            else:
                print(f"FAIL: Component 2 — Slide 1 notes mismatch")
                for i, (exp, act) in enumerate(zip(slide1_expected, slide1_actual)):
                    if exp != act:
                        print(f"  Line {i}: expected {repr(exp[:80])}, found {repr(act[:80])}")
        else:
            print(f"FAIL: Component 2 — Not enough paragraphs to check slide 1 content")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: All notes text content matches pptx notes in slide order (0.5 points)
    # Full comparison: every paragraph in docx must exactly match the corresponding
    # expected line derived from the pptx presenter notes.
    try:
        if doc_count == expected_count and expected_count > 0:
            mismatches = []
            for i, (exp, act) in enumerate(zip(expected_paras, doc_paras)):
                if exp != act:
                    mismatches.append((i, exp, act))

            if not mismatches:
                print(f"PASS: Component 3 — All {expected_count} paragraphs match pptx notes content (0.5 pts)")
                total_score += 0.5
            else:
                print(f"FAIL: Component 3 — {len(mismatches)} paragraph(s) do not match")
                # Show first 3 mismatches for debugging
                for idx, exp, act in mismatches[:3]:
                    print(f"  Para {idx}: expected {repr(exp[:80])}")
                    print(f"           actual   {repr(act[:80])}")
        elif doc_count != expected_count:
            print(f"FAIL: Component 3 — Cannot do full comparison due to paragraph count mismatch")
            # Partial credit: check what fraction of content is correct
            # Compare as many lines as possible (up to min length)
            min_len = min(doc_count, expected_count)
            if min_len > 0:
                match_count = sum(1 for e, a in zip(expected_paras[:min_len], doc_paras[:min_len]) if e == a)
                match_ratio = match_count / expected_count
                partial = round(match_ratio * 0.5, 2)
                if partial > 0:
                    print(f"  Partial: {match_count}/{expected_count} lines match, awarding {partial} pts")
                    total_score += partial
        else:
            print(f"FAIL: Component 3 — No expected content to compare against")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Run verification
verify_task(DOCX_PATH, PPTX_PATH)
