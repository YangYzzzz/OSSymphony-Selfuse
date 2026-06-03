"""
FINAL REWARD SCRIPT - SUCCESS
Task: Remove trailing spaces at the end of all paragraphs.
Generated: 2025-10-17 09:20:00
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation


def verify_trailing_spaces_removed(file_path: str) -> float:
    """
    Verify that there are **no trailing spaces or tabs** at the end of any
    paragraph across all text frames in the given PowerPoint presentation.

    Scoring (progressive):
        • For each non-empty paragraph:
            – 1 point if its last character is NOT whitespace
            – 0 points if it ends with whitespace (space, tab, etc.)
        • Final reward = (paragraphs_without_trailing_ws / total_non_empty_paragraphs)

    A perfect presentation (no trailing whitespace) yields 1.0.
    Partial deductions are applied proportionally otherwise.
    """

    print(f"Verifying trailing spaces in presentation: {file_path}")

    # Basic existence check – **no points awarded for existence**
    if not os.path.exists(file_path):
        print("✗ File not found – cannot verify task")
        return 0.0

    # Attempt to load the presentation – prerequisite, no score yet
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Failed to load PPTX: {e}")
        return 0.0

    total_paragraphs = 0      # Non-empty paragraphs we evaluate
    trailing_issues = 0       # Paragraphs ending with whitespace (failure cases)

    # Iterate through every paragraph in every text frame
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue  # Shapes without text are irrelevant

            for para_idx, para in enumerate(shape.text_frame.paragraphs, start=1):
                text = para.text or ""

                # Only score paragraphs that actually contain visible text
                if text == "":
                    continue

                total_paragraphs += 1

                # Detect trailing whitespace (space OR any other whitespace char)
                if text[-1].isspace():
                    trailing_issues += 1
                    snippet = text.replace("\n", " ")
                    print(
                        f"  ✗ Trailing whitespace found – Slide {slide_idx}, Paragraph {para_idx}:\n    › '{snippet}'"
                    )

    if total_paragraphs == 0:
        # No paragraphs to verify means task cannot be assessed (score 0)
        print("✗ No non-empty paragraphs found – nothing to evaluate")
        return 0.0

    paragraphs_without_issues = total_paragraphs - trailing_issues

    # Progressive scoring based on actual verification results
    score = paragraphs_without_issues / total_paragraphs

    print("\n–––– SUMMARY ––––")
    print(f"Total non-empty paragraphs           : {total_paragraphs}")
    print(f"Paragraphs WITHOUT trailing whitespace: {paragraphs_without_issues}")
    print(f"Paragraphs WITH trailing whitespace   : {trailing_issues}")
    print(f"Raw score                             : {score}")

    # Ensure score is a float within [0, 1]
    final_score = max(0.0, min(round(score, 2), 1.0))

    return final_score


if __name__ == "__main__":
    # Path to the presentation to verify – adjust if necessary
    FILE_PATH = "/home/user/remove_trailing_spaces_at_the_end_of_all_paragraphs.pptx"

    reward = verify_trailing_spaces_removed(FILE_PATH)
    print(f"REWARD: {reward}")

