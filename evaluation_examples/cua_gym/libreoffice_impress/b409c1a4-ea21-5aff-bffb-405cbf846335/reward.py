"""
FINAL REWARD SCRIPT - SUCCESS
Task: I'm updating our sales deck and want the chart on slide 7 to be a line chart instead of a bar chart, but I need to keep the original data the same. How do I make this change in Impress?
Generated: 2025-08-07 12:54:03
Status: success
Model: o4-mini
Total Steps: 8
"""

import os
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

file_path = '/home/user/im_updating_our_sales_deck_and_want_the_chart_on_slide_7_to_be_a_line_chart_instead_of_a_bar_chart_b.pptx'

def verify_task(fp):
    print("Checking task completion...")
    total_score = 0.0
    max_score = 1.0

    # 1. File existence and loading (0.2)
    if not os.path.exists(fp):
        print(f"✗ File not found: {fp}")
        print("Final score: 0.0")
        print("REWARD: 0.0")
        return 0.0
    print(f"✓ File exists: {fp}")
    try:
        prs = Presentation(fp)
        print(f"✓ Presentation loaded with {len(prs.slides)} slides")
        total_score += 0.2
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        print("Final score: 0.0")
        print("REWARD: 0.0")
        return 0.0

    # 2. Slide count >= 7 (0.2)
    slide_count = len(prs.slides)
    if slide_count >= 7:
        print(f"✓ Slide count {slide_count} >= 7")
        total_score += 0.2
    else:
        print(f"✗ Slide count {slide_count} < 7")

    # 3. Find chart on slide 7 (0.3)
    slide7 = prs.slides[6]
    chart_shape = None
    for shape in slide7.shapes:
        if hasattr(shape, 'chart'):
            chart_shape = shape
            break
    if chart_shape is None:
        print("✗ No chart found on slide 7")
    else:
        print("✓ Found chart on slide 7")
        total_score += 0.3
        chart = chart_shape.chart

        # 4. Chart type check (0.2)
        ctype = chart.chart_type
        if ctype in (XL_CHART_TYPE.LINE, XL_CHART_TYPE.LINE_MARKERS):
            print(f"✓ Chart type is line ({ctype})")
            total_score += 0.2
        else:
            print(f"✗ Chart type is not line (found {ctype})")

        # 5. Data preservation: series and categories check (0.1)
        try:
            plots = chart.plots
            if plots:
                plot = plots[0]
                series_count = len(list(plot.series))
                categories_count = len(list(plot.categories))
                print(f"Series count: {series_count}, Categories count: {categories_count}")
                if series_count > 0 and categories_count > 0:
                    print("✓ Chart data preserved (series and categories present)")
                    total_score += 0.1
                else:
                    print("✗ Chart data missing or empty")
            else:
                print("✗ No plots found in chart")
        except Exception as e:
            print(f"✗ Error checking chart data: {e}")

    # Final score computation and formatting
    final_score = min(total_score, max_score)
    if abs(final_score - max_score) < 1e-6:
        final_score = max_score
    final_score = round(final_score, 1)

    print(f"Final score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score

# Execute verification
enable = verify_task(file_path)
