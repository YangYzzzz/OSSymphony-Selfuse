"""
Initial Setup: Workshop Guide presentation with default bullets
Task ID: impress_ma_027
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_027'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Workshop Guide: Data Analytics Fundamentals"
    slide1.placeholders[1].text = "Q2 2025 Training Program\nTechVision Analytics Corp."

    # ---- Slide 2: Agenda (bulleted) ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Workshop Agenda"
    tf = slide2.placeholders[1].text_frame
    tf.text = "Introduction to Data Analytics"
    items_l1 = [
        "Data Collection Methods",
        "Statistical Analysis Techniques",
        "Data Visualization Best Practices",
        "Hands-On Lab Sessions",
        "Q&A and Wrap-Up",
    ]
    for item in items_l1:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

    # ---- Slide 3: Prerequisites (mixed levels) ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Prerequisites"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Software Requirements"
    items3 = [
        (1, "Python 3.10 or higher"),
        (1, "Jupyter Notebook"),
        (1, "Anaconda distribution recommended"),
        (0, "Knowledge Requirements"),
        (1, "Basic statistics understanding"),
        (1, "Familiarity with spreadsheets"),
        (0, "Hardware"),
        (1, "Laptop with 8GB RAM minimum"),
        (1, "Stable internet connection"),
    ]
    for level, text in items3:
        p = tf3.add_paragraph()
        p.text = text
        p.level = level

    # ---- Slide 4: Data Collection (bulleted) ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Module 1: Data Collection"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Primary Data Sources"
    items4 = [
        (1, "Surveys and questionnaires"),
        (1, "Direct observation studies"),
        (0, "Secondary Data Sources"),
        (1, "Government databases"),
        (1, "Industry reports from Gartner and McKinsey"),
        (1, "Academic research papers"),
        (0, "Data Quality Considerations"),
        (1, "Accuracy and completeness checks"),
        (1, "Timeliness of data collection"),
    ]
    for level, text in items4:
        p = tf4.add_paragraph()
        p.text = text
        p.level = level

    # ---- Slide 5: Statistical Analysis ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Module 2: Statistical Analysis"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Descriptive Statistics"
    items5 = [
        (1, "Mean, median, mode calculations"),
        (1, "Standard deviation and variance"),
        (0, "Inferential Statistics"),
        (1, "Hypothesis testing (t-test, chi-square)"),
        (1, "Confidence intervals"),
        (1, "Regression analysis"),
        (0, "Practical Applications"),
        (1, "Customer segmentation at RetailMax Inc."),
        (1, "Revenue forecasting for FY2025"),
    ]
    for level, text in items5:
        p = tf5.add_paragraph()
        p.text = text
        p.level = level

    # ---- Slide 6: Visualization ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Module 3: Data Visualization"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Chart Types and When to Use Them"
    items6 = [
        (1, "Bar charts for categorical comparisons"),
        (1, "Line charts for time series trends"),
        (1, "Scatter plots for correlation analysis"),
        (0, "Visualization Tools"),
        (1, "Matplotlib and Seaborn in Python"),
        (1, "Tableau Desktop and Tableau Public"),
        (1, "Power BI for enterprise dashboards"),
    ]
    for level, text in items6:
        p = tf6.add_paragraph()
        p.text = text
        p.level = level

    # ---- Slide 7: Lab Session 1 ----
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Lab Session 1: Exploratory Data Analysis"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "Dataset: GlobalSales_2024.csv"
    items7 = [
        (1, "12,500 transaction records"),
        (1, "15 feature columns"),
        (0, "Tasks"),
        (1, "Load and inspect the dataset"),
        (1, "Handle missing values (approx. 3% null rate)"),
        (1, "Generate summary statistics"),
        (1, "Create initial visualizations"),
        (0, "Deliverables"),
        (1, "Completed Jupyter notebook"),
        (1, "3-slide summary of findings"),
    ]
    for level, text in items7:
        p = tf7.add_paragraph()
        p.text = text
        p.level = level

    # ---- Slide 8: Lab Session 2 ----
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Lab Session 2: Predictive Modeling"
    tf8 = slide8.placeholders[1].text_frame
    tf8.text = "Building a Sales Forecast Model"
    items8 = [
        (1, "Feature engineering techniques"),
        (1, "Train/test split (80/20)"),
        (0, "Model Selection"),
        (1, "Linear regression baseline"),
        (1, "Random forest comparison"),
        (1, "XGBoost for advanced predictions"),
        (0, "Evaluation Metrics"),
        (1, "RMSE target: < $500"),
        (1, "R-squared target: > 0.85"),
    ]
    for level, text in items8:
        p = tf8.add_paragraph()
        p.text = text
        p.level = level

    # ---- Slide 9: Schedule ----
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Workshop Schedule"
    tf9 = slide9.placeholders[1].text_frame
    tf9.text = "Day 1 (March 25, 2025)"
    items9 = [
        (1, "09:00 - 10:30  Module 1: Data Collection"),
        (1, "10:45 - 12:15  Module 2: Statistical Analysis"),
        (1, "13:30 - 15:00  Lab Session 1"),
        (0, "Day 2 (March 26, 2025)"),
        (1, "09:00 - 10:30  Module 3: Data Visualization"),
        (1, "10:45 - 12:15  Lab Session 2"),
        (1, "13:30 - 15:00  Group Presentations"),
        (1, "15:15 - 16:00  Q&A and Certification"),
    ]
    for level, text in items9:
        p = tf9.add_paragraph()
        p.text = text
        p.level = level

    # ---- Slide 10: Contact / Closing ----
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "Contact & Resources"
    tf10 = slide10.placeholders[1].text_frame
    tf10.text = "Instructors"
    items10 = [
        (1, "Dr. Anika Patel - Lead Instructor"),
        (1, "James Rodriguez - Lab Facilitator"),
        (0, "Resources"),
        (1, "Course materials: learn.techvision.com/da-workshop"),
        (1, "Slack channel: #da-workshop-q2-2025"),
        (0, "Certification"),
        (1, "Complete both lab sessions"),
        (1, "Score 80%+ on assessment quiz"),
        (1, "Certificate issued within 5 business days"),
    ]
    for level, text in items10:
        p = tf10.add_paragraph()
        p.text = text
        p.level = level

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
