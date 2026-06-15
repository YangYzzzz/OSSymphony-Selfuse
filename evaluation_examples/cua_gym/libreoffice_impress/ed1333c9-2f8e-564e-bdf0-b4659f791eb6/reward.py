"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 38 I’d like to swap the standard bullet symbol for the “Dash” style and recolor those bullets to the palette’s exact Blue 7 shade. How do I do that in LibreOffice Impress?
Generated: 2025-09-10 22:06:59
Status: success
Model: azure-o3
Total Steps: 11
"""

import os
from pptx import Presentation
from pptx.dml.color import RGBColor

"""
Reward Script: verify_slide_38_bullets

Task to verify:
On slide 38 the standard bullet symbol should be swapped for the “Dash” style and
those bullets must be recolored to the palette’s exact Blue 7 shade (#2A6099).

Scoring (progressive):
• 0.6 – All bullet paragraphs on slide 38 use a dash character (–, -, —)
• 0.4 – The dash bullets are colored exactly Blue 7 (RGB 2A6099)
The script awards partial credit if only one of the objectives is satisfied.

Verification methodology:
1. Load the PPTX with python-pptx (file existence & load get NO points).
2. Ensure slide 38 exists (index 37).
3. Iterate through every text-bearing shape on slide 38, locating paragraphs that
   look like bullets (first non-space char is either a dash or a “•”).
4. For each bullet paragraph found:
      • Confirm the first visible character is a dash → dash requirement.
      • Inspect the first run’s font color and compare to RGB 2A6099 → color req.
5. Aggregate results and return a score between 0.0 and 1.0.
6. Print detailed diagnostics plus final line:  "REWARD: X.X".
"""

target_rgb = RGBColor(0x2A, 0x60, 0x99)  # Blue 7 hex


def _is_dash(char: str) -> bool:
    """Return True if char is an accepted dash bullet character."""
    return char in {"–", "-", "—"}


def verify_slide_38_bullets(file_path: str) -> float:
    print(f"Loading presentation: {file_path}")

    # Preliminary checks (no points for these!)
    if not os.path.exists(file_path):
        print("✗ File not found")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Could not load presentation: {e}")
        return 0.0

    # Slide index validation
    if len(prs.slides) < 38:
        print(f"✗ Only {len(prs.slides)} slides present – need at least 38")
        return 0.0
    slide = prs.slides[37]
    print("✓ Slide 38 located")

    dash_ok = True   # assume true until disproved
    color_ok = True
    bullet_paragraphs = 0

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for p_idx, para in enumerate(shape.text_frame.paragraphs, start=1):
            text = para.text.strip()
            if not text:
                continue

            first_char = text[0]
            if _is_dash(first_char) or first_char == "•":  # bullet-looking para
                bullet_paragraphs += 1

                # Check dash replacement
                if _is_dash(first_char):
                    print(f"✓ Paragraph {p_idx}: uses dash bullet")
                else:
                    dash_ok = False
                    print(f"✗ Paragraph {p_idx}: bullet is '{first_char}', not dash")

                # Check color of dash (first run normally contains the bullet char)
                if para.runs:
                    first_run = para.runs[0]
                    rgb = getattr(first_run.font.color, "rgb", None)
                    if rgb == target_rgb:
                        print(f"✓ Paragraph {p_idx}: bullet color matches Blue 7 ({rgb})")
                    else:
                        color_ok = False
                        print(f"✗ Paragraph {p_idx}: bullet color {rgb} ≠ Blue 7 {target_rgb}")
                else:
                    color_ok = False
                    print(f"✗ Paragraph {p_idx}: no runs detected – cannot verify color")

    # If no bullet paragraphs found, the task clearly fails.
    if bullet_paragraphs == 0:
        print("✗ No bullet paragraphs detected on slide 38")
        return 0.0

    # Progressive scoring
    score = 0.0
    if dash_ok:
        score += 0.6
    if color_ok:
        score += 0.4

    print(f"Total score: {score}")
    return score


if __name__ == "__main__":
    FILE = "/home/user/on_slide_38_id_like_to_swap_the_standard_bullet_symbol_for_the_dash_style_and_recolor_those_bullets__golden.pptx"
    reward = verify_slide_38_bullets(FILE)
    print(f"REWARD: {reward}")
