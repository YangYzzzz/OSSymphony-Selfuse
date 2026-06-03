"""
Initial Setup: Configure master slide font scheme presentation
Task ID: impress_gf2_043
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_043'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_placeholder_font(placeholder, font_name, font_size_pt=None):
    """Set font on all runs in a placeholder's text frame, and set default run props."""
    if not placeholder.has_text_frame:
        return
    tf = placeholder.text_frame
    # Set default paragraph run properties so new text inherits this font
    for para in tf.paragraphs:
        # Set on existing runs
        for run in para.runs:
            run.font.name = font_name
            if font_size_pt:
                run.font.size = Pt(font_size_pt)
        # Set default run properties for the paragraph
        pPr = para._p.get_or_add_pPr()
        defRPr = pPr.find(qn('a:defRPr'))
        if defRPr is None:
            defRPr = pPr.makeelement(qn('a:defRPr'), {})
            pPr.append(defRPr)
        # Set latin font
        latin = defRPr.find(qn('a:latin'))
        if latin is None:
            latin = defRPr.makeelement(qn('a:latin'), {})
            defRPr.append(latin)
        latin.set('typeface', font_name)
        if font_size_pt:
            defRPr.set('sz', str(int(font_size_pt * 100)))  # hundredths of a point


def create_initial():
    prs = Presentation()

    # Access the slide master
    slide_master = prs.slide_masters[0]

    # Set ALL master slide placeholders to 'Liberation Sans' (the initial default font)
    for ph in slide_master.placeholders:
        set_placeholder_font(ph, 'Liberation Sans')

    # Also set font on all slide layouts' placeholders
    for layout in prs.slide_layouts:
        for ph in layout.placeholders:
            set_placeholder_font(ph, 'Liberation Sans')

    # Create 16 slides with varied content
    slide_data = [
        (0, "Formal Report", "Q1 2025 Performance Review"),
        (1, "Executive Summary", [
            "Revenue grew 12% year-over-year to $4.2M",
            "Customer acquisition cost decreased by 8%",
            "Net promoter score improved from 72 to 81",
            "Three new product lines launched successfully",
        ]),
        (1, "Market Analysis", [
            "Total addressable market expanded to $12.8B",
            "Competitive landscape shifted with two new entrants",
            "Customer segments showed varying growth patterns",
            "Digital transformation accelerated across industries",
        ]),
        (1, "Financial Highlights", [
            "Gross margin improved to 68.3% from 65.1%",
            "Operating expenses reduced by $340K",
            "Cash reserves stand at $8.7M",
            "Accounts receivable turnover improved to 45 days",
        ]),
        (1, "Product Development", [
            "Platform 3.0 released with 47 new features",
            "Mobile app downloads exceeded 250,000",
            "API integration partnerships grew to 23",
            "User engagement metrics up 31% quarter-over-quarter",
        ]),
        (1, "Customer Success", [
            "Enterprise client retention at 96.4%",
            "Average contract value increased to $48,500",
            "Support ticket resolution time down to 2.3 hours",
            "Launched dedicated success manager program",
        ]),
        (5, None, None),  # Blank slide - section divider
        (1, "Sales Performance", [
            "North America: $2.1M (+15% YoY)",
            "Europe: $1.3M (+9% YoY)",
            "Asia-Pacific: $580K (+22% YoY)",
            "Latin America: $220K (+34% YoY)",
        ]),
        (1, "Operations Update", [
            "Server uptime maintained at 99.97%",
            "Infrastructure costs optimized by 18%",
            "Deployed to three new data center regions",
            "Security audit completed with zero critical findings",
        ]),
        (1, "Human Resources", [
            "Headcount grew from 124 to 148 employees",
            "Employee satisfaction score: 4.3 out of 5.0",
            "Voluntary turnover reduced to 8.2%",
            "Launched mentorship and leadership programs",
        ]),
        (1, "Marketing Initiatives", [
            "Brand awareness increased 25% in target demographics",
            "Content marketing generated 12,000 qualified leads",
            "Social media following grew to 85,000 across platforms",
            "Event sponsorships yielded 340 enterprise prospects",
        ]),
        (1, "Technology Roadmap", [
            "AI-powered analytics engine in beta testing",
            "Microservices migration 73% complete",
            "Real-time collaboration features planned for Q3",
            "Enhanced security framework deployment in Q2",
        ]),
        (1, "Risk Assessment", [
            "Supply chain dependencies monitored and diversified",
            "Regulatory compliance updated for new data privacy laws",
            "Cybersecurity insurance coverage expanded",
            "Business continuity plans tested and validated",
        ]),
        (1, "Sustainability Goals", [
            "Carbon footprint reduced by 22% through cloud optimization",
            "Remote work policy decreased office energy consumption",
            "Partnered with two environmental nonprofits",
            "ESG reporting framework adopted for transparency",
        ]),
        (1, "Q2 Priorities", [
            "Expand enterprise sales team by 6 representatives",
            "Launch Platform 3.1 with advanced reporting",
            "Enter Japanese market with localized product",
            "Achieve SOC 2 Type II certification",
        ]),
        (1, "Thank You", [
            "Questions and discussion welcome",
            "Contact: leadership@formalreport.example.com",
            "Next review scheduled: July 15, 2025",
        ]),
    ]

    for layout_idx, title_text, content in slide_data:
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        if title_text and slide.shapes.title:
            slide.shapes.title.text = title_text
            # Ensure title font is Liberation Sans
            for para in slide.shapes.title.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = 'Liberation Sans'

        if content and layout_idx == 1:
            # Content placeholder (index 1)
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            if isinstance(content, list):
                for i, item in enumerate(content):
                    if i == 0:
                        tf.paragraphs[0].text = item
                        for run in tf.paragraphs[0].runs:
                            run.font.name = 'Liberation Sans'
                    else:
                        p = tf.add_paragraph()
                        p.text = item
                        for run in p.runs:
                            run.font.name = 'Liberation Sans'
            else:
                tf.paragraphs[0].text = str(content)
                for run in tf.paragraphs[0].runs:
                    run.font.name = 'Liberation Sans'

        if layout_idx == 0 and isinstance(content, str):
            # Subtitle for title slide
            if 1 in slide.placeholders:
                slide.placeholders[1].text = content
                for para in slide.placeholders[1].text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = 'Liberation Sans'

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
