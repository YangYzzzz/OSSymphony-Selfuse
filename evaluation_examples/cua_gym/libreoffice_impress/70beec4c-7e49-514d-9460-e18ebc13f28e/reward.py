"""
FINAL REWARD SCRIPT - SUCCESS
Task: Set Table 1 column widths to 2.0 cm, 3.5 cm, 3.5 cm, 2.0 cm.
Generated: 2025-10-17 07:37:34
Status: success
Model: azure-o3
Total Steps: 1
"""

from pptx import Presentation
import os

# Constants
EMU_PER_CM = 360000        # PowerPoint stores lengths in English Metric Units
TOLERANCE_EMU = 10000      # ≈ 0.028 cm – small leniency for rounding

# Expected widths for the first four columns (Table 1) in cm
EXPECTED_WIDTHS_CM = [2.0, 3.5, 3.5, 2.0]
EXPECTED_WIDTHS_EMU = [int(round(cm * EMU_PER_CM)) for cm in EXPECTED_WIDTHS_CM]

FILE_PATH = "/home/user/set_table_1_column_widths_to_20_cm_35_cm_35_cm_20_cm.pptx"


def verify_table_1_column_widths(file_path: str) -> float:
    """Verify that the first table's four column widths match the specification.

    Returns a progressive score between 0.0 and 1.0.
    """
    # 0.  Basic file presence check (no points awarded)
    if not os.path.exists(file_path):
        print("✗ File not found:", file_path)
        return 0.0

    # 1.  Load the presentation (prerequisite – no points yet)
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        return 0.0

    # 2.  Locate the first table in the presentation (Table 1)
    table_shape = None
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_table:
                table_shape = shape
                print(f"✓ Found table on slide {slide_idx + 1}, shape {shape_idx + 1}")
                break
        if table_shape:
            break

    if table_shape is None:
        print("✗ No table found in the presentation")
        return 0.0

    table = table_shape.table
    num_columns = len(table.columns)
    print(f"Table contains {num_columns} columns")

    # 3.  Verify the widths of the first four columns
    correct_columns = 0
    for idx in range(4):
        if idx >= num_columns:
            print(f"✗ Table has fewer than {idx + 1} columns – cannot verify further")
            break

        actual_width = table.columns[idx].width
        expected_width = EXPECTED_WIDTHS_EMU[idx]
        diff = abs(actual_width - expected_width)
        actual_width_cm = actual_width / EMU_PER_CM

        print(
            f"Column {idx + 1}: actual = {actual_width_cm:.2f} cm, "
            f"expected = {EXPECTED_WIDTHS_CM[idx]} cm, diff = {diff} EMU"
        )

        if diff <= TOLERANCE_EMU:
            print("  ✓ Width within tolerance")
            correct_columns += 1
        else:
            print("  ✗ Width outside tolerance")

    # 4.  Progressive scoring – each correctly-sized column earns 0.25
    score = correct_columns / 4.0
    print(f"Total correct columns: {correct_columns}/4")
    print(f"REWARD: {score}")
    return score


# ---------------------------------------------------------------------------
# Execute verification when script is run
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    verify_table_1_column_widths(FILE_PATH)
