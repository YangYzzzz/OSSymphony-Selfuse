"""
Initial Setup: Insert a bar chart on slide 2 showing quarterly revenue
Task ID: impress_tct_036
Domain: libreoffice_impress

Creates a 5-slide presentation. Slide 2 has title 'Quarterly Revenue' but NO chart.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_036'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Annual Revenue Report"
    slide1.placeholders[1].text = "FY 2025 Financial Overview\nPrepared by Finance Department"

    # --- Slide 2: Quarterly Revenue (NO chart - agent task is to add one) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title text box
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Quarterly Revenue"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x6E)

    # Add a subtitle text
    txBox2 = slide2.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Revenue breakdown by quarter for fiscal year 2025 (in thousands USD)"
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.name = "Arial"
    run2.font.size = Pt(14)
    run2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # --- Slide 3: Regional Breakdown ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Regional Performance"
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.runs[0]
    r3.font.name = "Arial"
    r3.font.size = Pt(28)
    r3.font.bold = True
    r3.font.color.rgb = RGBColor(0x2E, 0x4A, 0x6E)

    # Add a table with regional data
    table_shape = slide3.shapes.add_table(5, 3, Inches(1.5), Inches(1.8), Inches(7), Inches(2.5))
    table = table_shape.table
    headers = ["Region", "Revenue (K)", "Growth %"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)

    regions = [
        ["North America", "$345K", "+12.3%"],
        ["Europe", "$218K", "+8.7%"],
        ["Asia Pacific", "$127K", "+15.1%"],
        ["Latin America", "$65K", "+4.2%"],
    ]
    for r_idx, row_data in enumerate(regions, 1):
        for c_idx, val in enumerate(row_data):
            table.cell(r_idx, c_idx).text = val

    # --- Slide 4: Key Highlights ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Key Highlights"
    p4.alignment = PP_ALIGN.CENTER
    r4 = p4.runs[0]
    r4.font.name = "Arial"
    r4.font.size = Pt(28)
    r4.font.bold = True
    r4.font.color.rgb = RGBColor(0x2E, 0x4A, 0x6E)

    bullets_box = slide4.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
    btf = bullets_box.text_frame
    btf.word_wrap = True

    highlights = [
        "Total annual revenue reached $690K, exceeding target by 7.4%",
        "Q3 posted the highest quarterly revenue at $210K",
        "Customer acquisition cost decreased by 14% year-over-year",
        "Enterprise segment grew 22%, driven by new partnerships",
        "Recurring revenue now accounts for 68% of total revenue",
    ]
    for i, text in enumerate(highlights):
        if i == 0:
            p = btf.paragraphs[0]
        else:
            p = btf.add_paragraph()
        p.text = text
        p.level = 0
        p.space_after = Pt(12)
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(16)

    # --- Slide 5: Next Steps ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf5 = txBox5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Next Steps & Outlook"
    p5.alignment = PP_ALIGN.CENTER
    r5 = p5.runs[0]
    r5.font.name = "Arial"
    r5.font.size = Pt(28)
    r5.font.bold = True
    r5.font.color.rgb = RGBColor(0x2E, 0x4A, 0x6E)

    outlook_box = slide5.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
    otf = outlook_box.text_frame
    otf.word_wrap = True

    steps = [
        "Expand enterprise sales team in Asia Pacific region",
        "Launch premium subscription tier in Q2 2026",
        "Invest in AI-driven product features to improve retention",
        "Target 15% overall revenue growth for FY 2026",
    ]
    for i, text in enumerate(steps):
        if i == 0:
            p = otf.paragraphs[0]
        else:
            p = otf.add_paragraph()
        p.text = text
        p.level = 0
        p.space_after = Pt(12)
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(16)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
