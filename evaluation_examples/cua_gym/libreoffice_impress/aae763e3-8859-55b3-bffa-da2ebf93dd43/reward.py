"""
Reward Script: Insert a bar chart on slide 2 showing quarterly revenue
Task ID: impress_tct_036
Domain: libreoffice_impress
Scoring:
  Component 1: Chart exists on slide 2 (0.3 pts)
  Component 2: Chart type is bar/column clustered (0.2 pts)
  Component 3: Chart data values match [120, 185, 210, 175] (0.3 pts)
  Component 4: Chart categories match ['Q1', 'Q2', 'Q3', 'Q4'] (0.2 pts)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_036'


def persist_app_state(domain: str):
    """Save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.chart import XL_CHART_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Verify we have at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 2")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]  # 0-indexed, slide 2

    # Find chart shapes on slide 2
    chart_shapes = [s for s in slide2.shapes if s.has_chart]

    # Component 1: Chart exists on slide 2 (0.3 points)
    try:
        if len(chart_shapes) > 0:
            print(f"PASS: Component 1 -- Chart found on slide 2 ({len(chart_shapes)} chart(s)) (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 -- No chart found on slide 2")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    if len(chart_shapes) == 0:
        # No chart means remaining components cannot pass
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    chart = chart_shapes[0].chart

    # Component 2: Chart type is bar/column (0.2 points)
    # python-pptx uses COLUMN_CLUSTERED (51) for vertical bar charts
    # Also accept BAR_CLUSTERED (57) as horizontal bar charts
    try:
        chart_type = chart.chart_type
        # Column clustered (vertical bars) and bar clustered (horizontal bars)
        # are both valid "bar charts"
        valid_bar_types = [
            XL_CHART_TYPE.COLUMN_CLUSTERED,      # 51
            XL_CHART_TYPE.BAR_CLUSTERED,          # 57
            XL_CHART_TYPE.COLUMN_STACKED,         # 52
            XL_CHART_TYPE.BAR_STACKED,            # 58
        ]
        if chart_type in valid_bar_types:
            print(f"PASS: Component 2 -- Chart type is bar/column: {chart_type} (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 2 -- Expected bar/column chart, found type: {chart_type}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Chart data values match [120, 185, 210, 175] (0.3 points)
    try:
        expected_values = [120.0, 185.0, 210.0, 175.0]
        if len(chart.series) > 0:
            actual_values = list(chart.series[0].values)
            # Compare with tolerance for floating point
            values_match = (len(actual_values) == len(expected_values) and
                           all(abs(a - e) < 0.01 for a, e in zip(actual_values, expected_values)))
            if values_match:
                print(f"PASS: Component 3 -- Data values match: {actual_values} (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 3 -- Expected values {expected_values}, found {actual_values}")
        else:
            print(f"FAIL: Component 3 -- No data series found in chart")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Chart categories match ['Q1', 'Q2', 'Q3', 'Q4'] (0.2 points)
    try:
        expected_cats = ['Q1', 'Q2', 'Q3', 'Q4']
        actual_cats = [str(c) for c in chart.plots[0].categories]
        if actual_cats == expected_cats:
            print(f"PASS: Component 4 -- Categories match: {actual_cats} (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 4 -- Expected categories {expected_cats}, found {actual_cats}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
