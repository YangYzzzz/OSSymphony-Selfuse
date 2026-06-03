"""
Initial Setup: Create a 7-slide Molecular Biology presentation with slide 4 empty.
Task ID: impress_teach_027
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_027'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_content_slide(prs, layout_idx, title_text, body_lines):
    """Add a slide with a title and bullet-point body text."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title_text
    # Find the body placeholder (index 1)
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Molecular Biology"
    slide1.placeholders[1].text = "Fundamentals of Life at the Molecular Level"

    # --- Slide 2: Introduction to Biomolecules ---
    add_content_slide(prs, 1, "Introduction to Biomolecules", [
        "Living organisms are composed of four major biomolecules",
        "Nucleic acids (DNA and RNA) store genetic information",
        "Proteins perform most cellular functions",
        "Carbohydrates provide energy and structural support",
        "Lipids form cell membranes and store energy",
    ])

    # --- Slide 3: DNA Structure ---
    add_content_slide(prs, 1, "DNA Structure", [
        "Double helix composed of two antiparallel strands",
        "Nucleotides contain deoxyribose sugar, phosphate, and a nitrogenous base",
        "Four bases: Adenine (A), Thymine (T), Guanine (G), Cytosine (C)",
        "Base pairing rules: A-T (2 hydrogen bonds), G-C (3 hydrogen bonds)",
        "5' to 3' directionality is critical for replication",
    ])

    # --- Slide 4: Comparison of Biomolecules (EMPTY - no table) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add only a title text box
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Comparison of Biomolecules"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # --- Slide 5: RNA Types ---
    add_content_slide(prs, 1, "RNA Types", [
        "Messenger RNA (mRNA) carries genetic code from DNA to ribosomes",
        "Transfer RNA (tRNA) delivers amino acids during translation",
        "Ribosomal RNA (rRNA) forms the structural core of ribosomes",
        "MicroRNA (miRNA) regulates gene expression post-transcriptionally",
        "Small nuclear RNA (snRNA) involved in RNA splicing",
    ])

    # --- Slide 6: Protein Functions ---
    add_content_slide(prs, 1, "Protein Functions", [
        "Enzymes catalyze biochemical reactions (e.g., DNA polymerase)",
        "Structural proteins provide support (e.g., collagen, keratin)",
        "Transport proteins carry molecules (e.g., hemoglobin, albumin)",
        "Signaling proteins transmit messages (e.g., insulin, growth factors)",
        "Defense proteins protect organisms (e.g., antibodies, complement)",
    ])

    # --- Slide 7: Summary ---
    add_content_slide(prs, 1, "Summary", [
        "DNA stores the genetic blueprint in a stable double helix",
        "RNA serves as an intermediary for gene expression",
        "Proteins execute diverse cellular functions",
        "Understanding molecular biology is key to medicine and biotechnology",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
