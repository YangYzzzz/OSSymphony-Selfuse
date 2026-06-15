"""
Reward Script: Create a comparison table on slide 4 with specific headers and formatting
Task ID: impress_teach_027
Domain: libreoffice_impress
Scoring:
  Component 1: Table exists on slide 4 with correct dimensions (5 rows x 4 cols) — 0.25 pts
  Component 2: Header text matches ['Feature', 'DNA', 'RNA', 'Protein'] — 0.25 pts
  Component 3: Header cells have dark teal (#004D40) background fill — 0.25 pts
  Component 4: Header text is white (#FFFFFF) and bold — 0.25 pts
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_027'


def persist_app_state(domain):
    """Save any unsaved changes in LibreOffice before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.oxml.ns import qn
    except ImportError as e:
        print(f"CRITICAL: Missing library: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide4 = prs.slides[3]  # 0-indexed

    # Find the table on slide 4
    table = None
    for shape in slide4.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            break

    # Component 1: Table exists with correct dimensions (5 rows x 4 cols) — 0.25 pts
    try:
        if table is None:
            print("FAIL: Component 1 — No table found on slide 4")
        else:
            n_rows = len(table.rows)
            n_cols = len(table.columns)
            if n_rows == 5 and n_cols == 4:
                print(f"PASS: Component 1 — Table found with {n_rows} rows x {n_cols} cols (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 — Table is {n_rows}x{n_cols}, expected 5x4")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # If no table, remaining checks cannot proceed
    if table is None:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: Header text matches ['Feature', 'DNA', 'RNA', 'Protein'] — 0.25 pts
    expected_headers = ['Feature', 'DNA', 'RNA', 'Protein']
    try:
        actual_headers = []
        for c in range(min(len(table.columns), 4)):
            actual_headers.append(table.cell(0, c).text.strip())
        if actual_headers == expected_headers:
            print(f"PASS: Component 2 — Headers match: {actual_headers} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 — Headers are {actual_headers}, expected {expected_headers}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Header cells have dark teal (#004D40) background fill — 0.25 pts
    try:
        fill_ok_count = 0
        for c in range(min(len(table.columns), 4)):
            cell = table.cell(0, c)
            tcPr = cell._tc.tcPr
            fill_color = None
            if tcPr is not None:
                solidFill = tcPr.find(qn('a:solidFill'))
                if solidFill is not None:
                    srgb = solidFill.find(qn('a:srgbClr'))
                    if srgb is not None:
                        fill_color = srgb.get('val')
            if fill_color and fill_color.upper() == '004D40':
                fill_ok_count += 1
            else:
                print(f"  Cell(0,{c}) fill: {fill_color} (expected 004D40)")

        if fill_ok_count == 4:
            print(f"PASS: Component 3 — All 4 header cells have #004D40 background (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 — Only {fill_ok_count}/4 header cells have #004D40 fill")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Header text is white (#FFFFFF) and bold — 0.25 pts
    try:
        fmt_ok_count = 0
        for c in range(min(len(table.columns), 4)):
            cell = table.cell(0, c)
            cell_ok = False
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    is_bold = run.font.bold is True
                    is_white = False
                    try:
                        if run.font.color.type is not None:
                            is_white = str(run.font.color.rgb).upper() == 'FFFFFF'
                    except Exception:
                        pass
                    if is_bold and is_white:
                        cell_ok = True
                    else:
                        print(f"  Cell(0,{c}) run '{run.text}': bold={run.font.bold}, color={'white' if is_white else 'not white'}")
            if cell_ok:
                fmt_ok_count += 1

        if fmt_ok_count == 4:
            print(f"PASS: Component 4 — All 4 header cells have white bold text (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 4 — Only {fmt_ok_count}/4 header cells have white bold text")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
