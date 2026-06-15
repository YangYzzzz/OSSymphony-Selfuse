"""
Initial Setup: Build a cause-and-effect diagram presentation
Task ID: impress_teach_064
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_064'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Environmental Science"
    slide1.placeholders[1].text = "A Comprehensive Overview\nSpring 2025 Lecture Series"

    # --- Slide 2: Global Warming Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Global Warming Overview"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Average global temperature has risen by 1.1\u00b0C since pre-industrial era"
    p2a = tf2.add_paragraph()
    p2a.text = "2023 recorded as the hottest year in 125,000 years"
    p2a.level = 0
    p2b = tf2.add_paragraph()
    p2b.text = "Arctic sea ice declining at 13% per decade since 1979"
    p2b.level = 0
    p2c = tf2.add_paragraph()
    p2c.text = "Sea levels have risen approximately 20 cm over the past century"
    p2c.level = 0

    # --- Slide 3: Ocean Acidification ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Ocean Acidification"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Ocean pH has decreased by 0.1 units since the Industrial Revolution"
    p3a = tf3.add_paragraph()
    p3a.text = "Coral reefs face bleaching events with increasing frequency"
    p3a.level = 0
    p3b = tf3.add_paragraph()
    p3b.text = "Marine food chains disrupted by shell-forming organism decline"
    p3b.level = 0
    p3c = tf3.add_paragraph()
    p3c.text = "Projected pH drop of 0.3-0.4 by 2100 under current emission trajectories"
    p3c.level = 0

    # --- Slide 4: Biodiversity Loss ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Biodiversity Loss"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Over 1 million species currently threatened with extinction"
    p4a = tf4.add_paragraph()
    p4a.text = "Wildlife populations have declined by 69% since 1970 (WWF Living Planet Report)"
    p4a.level = 0
    p4b = tf4.add_paragraph()
    p4b.text = "Deforestation destroys approximately 10 million hectares annually"
    p4b.level = 0
    p4c = tf4.add_paragraph()
    p4c.text = "Pollinator species decline threatens global food security"
    p4c.level = 0

    # --- Slide 5: Causes of Climate Change (EMPTY - task target) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add only a title text box
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.0))
    tf5 = txBox.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Causes of Climate Change"
    p5.alignment = PP_ALIGN.CENTER
    run5 = p5.runs[0]
    run5.font.size = Pt(36)
    run5.font.bold = True
    run5.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 6: Mitigation Strategies ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Mitigation Strategies"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Transition to renewable energy sources (solar, wind, hydroelectric)"
    p6a = tf6.add_paragraph()
    p6a.text = "Implement carbon capture and storage technologies at scale"
    p6a.level = 0
    p6b = tf6.add_paragraph()
    p6b.text = "Strengthen international climate agreements and carbon pricing mechanisms"
    p6b.level = 0
    p6c = tf6.add_paragraph()
    p6c.text = "Promote sustainable agriculture and reforestation programs"
    p6c.level = 0

    # --- Slide 7: Conclusion ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Conclusion & Call to Action"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "Climate change is the defining challenge of our generation"
    p7a = tf7.add_paragraph()
    p7a.text = "Every fraction of a degree matters for ecosystems and human communities"
    p7a.level = 0
    p7b = tf7.add_paragraph()
    p7b.text = "Individual actions combined with systemic policy change can make a difference"
    p7b.level = 0
    p7c = tf7.add_paragraph()
    p7c.text = "The next decade is critical for achieving net-zero emission targets"
    p7c.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
