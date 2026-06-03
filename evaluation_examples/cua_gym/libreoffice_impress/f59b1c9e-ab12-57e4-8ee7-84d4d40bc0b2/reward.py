"""
Reward Script: Cause-and-effect diagram on slide 5
Task ID: impress_teach_064
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Right arrow shape labeled 'Climate Change' with fill #E53935
  Component 2 (0.35): Four rectangles labeled CO2 Emissions, Deforestation, Industrial Waste, Agriculture
  Component 3 (0.20): Four diagonal line connectors present on slide 5
  Component 4 (0.15): Positional layout - arrow on right, rectangles on left
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_064'


def persist_app_state(domain: str):
    """Save any unsaved GUI edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify cause-and-effect diagram on slide 5.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # Slide 5 (0-indexed)

    # Collect shapes by type for analysis
    auto_shapes = []
    line_shapes = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            auto_shapes.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
            line_shapes.append(shape)

    # Component 1: Right arrow shape labeled 'Climate Change' with fill #E53935 (0.30 points)
    # The shape must have SOLID fill #E53935 AND text matching 'Climate Change' (not a longer string
    # like 'Causes of Climate Change' which is the pre-existing slide title).
    try:
        arrow_found = False
        arrow_fill_correct = False
        arrow_text_correct = False

        for shape in auto_shapes:
            text = shape.text.strip() if hasattr(shape, 'text') and shape.text else ""
            text_lower = text.lower()

            # First gate: must have solid fill with #E53935 to be the arrow shape
            # (this distinguishes from pre-existing textboxes)
            has_target_fill = False
            try:
                fill = shape.fill
                if fill.type is not None and fill.type == 1:  # SOLID
                    rgb = str(fill.fore_color.rgb).upper()
                    if rgb == 'E53935':
                        has_target_fill = True
            except:
                pass

            if not has_target_fill:
                continue  # Skip shapes without the target fill color

            arrow_fill_correct = True

            # Check text contains 'climate change'
            if 'climate change' in text_lower:
                arrow_text_correct = True

            # Check if it's an arrow-type shape
            try:
                auto_type = shape.auto_shape_type
                if auto_type is not None and 'ARROW' in str(auto_type):
                    arrow_found = True
            except:
                pass

        if arrow_found and arrow_text_correct and arrow_fill_correct:
            print(f"PASS: Component 1 -- Right arrow 'Climate Change' with fill #E53935 (0.30 pts)")
            total_score += 0.30
        elif arrow_text_correct and arrow_fill_correct:
            # Has the text and color but may not be strictly arrow type
            print(f"PASS: Component 1 -- Shape 'Climate Change' with fill #E53935 found (0.25 pts)")
            total_score += 0.25
        elif arrow_fill_correct:
            # Has the fill color but text is wrong/missing
            print(f"PARTIAL: Component 1 -- Shape with fill #E53935 found but text missing/wrong (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 -- No shape with fill #E53935 found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Four rectangles with correct labels (0.35 points)
    # Each label is worth 0.35/4 = 0.0875 points
    try:
        expected_labels = ['co2 emissions', 'deforestation', 'industrial waste', 'agriculture']
        found_labels = set()

        for shape in auto_shapes:
            text = shape.text.strip().lower() if hasattr(shape, 'text') and shape.text else ""
            for label in expected_labels:
                if label in text:
                    # Verify it's a rectangle-ish shape (not the arrow)
                    try:
                        auto_type = shape.auto_shape_type
                        if auto_type is not None and 'ARROW' in str(auto_type):
                            continue  # Skip arrow shapes
                    except:
                        pass
                    found_labels.add(label)

        label_count = len(found_labels)
        label_score = (label_count / 4.0) * 0.35

        if label_count == 4:
            print(f"PASS: Component 2 -- All 4 rectangle labels found: {found_labels} (0.35 pts)")
        elif label_count > 0:
            print(f"PARTIAL: Component 2 -- {label_count}/4 rectangle labels found: {found_labels} ({label_score:.3f} pts)")
        else:
            print(f"FAIL: Component 2 -- No expected rectangle labels found on slide 5")

        total_score += label_score
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Four diagonal line connectors present (0.20 points)
    # Each connector is worth 0.05 points
    try:
        connector_count = len(line_shapes)

        if connector_count >= 4:
            print(f"PASS: Component 3 -- {connector_count} line connectors found (need >=4) (0.20 pts)")
            total_score += 0.20
        elif connector_count > 0:
            conn_score = (min(connector_count, 4) / 4.0) * 0.20
            print(f"PARTIAL: Component 3 -- {connector_count}/4 line connectors found ({conn_score:.3f} pts)")
            total_score += conn_score
        else:
            print(f"FAIL: Component 3 -- No line connectors found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Positional layout - arrow on right, rectangles on left (0.15 points)
    try:
        # Find the arrow shape (identified by #E53935 fill) and rectangles
        arrow_left_pos = None
        rect_positions = []

        for shape in auto_shapes:
            text = shape.text.strip().lower() if hasattr(shape, 'text') and shape.text else ""

            # Identify the arrow by its fill color
            has_target_fill = False
            try:
                fill = shape.fill
                if fill.type is not None and fill.type == 1:
                    rgb = str(fill.fore_color.rgb).upper()
                    if rgb == 'E53935':
                        has_target_fill = True
            except:
                pass

            if has_target_fill and 'climate change' in text:
                arrow_left_pos = shape.left
            elif text in ['co2 emissions', 'deforestation', 'industrial waste', 'agriculture']:
                rect_positions.append(shape.left)

        if arrow_left_pos is not None and len(rect_positions) > 0:
            # Arrow should be to the right of all rectangles
            all_rects_left_of_arrow = all(rp < arrow_left_pos for rp in rect_positions)
            if all_rects_left_of_arrow:
                print(f"PASS: Component 4 -- Arrow on right, rectangles on left (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 4 -- Layout incorrect: arrow_left={arrow_left_pos}, rect_lefts={rect_positions}")
        else:
            print(f"FAIL: Component 4 -- Cannot verify layout: arrow_pos={arrow_left_pos}, rect_count={len(rect_positions)}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.3f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
