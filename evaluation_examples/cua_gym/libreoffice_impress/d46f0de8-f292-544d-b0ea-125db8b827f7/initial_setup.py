"""
Initial Setup: Market Sizing presentation with empty funnel slide
Task ID: impress_exec_062
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_062'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_text(slide, text):
    """Set the title placeholder text."""
    if slide.shapes.title:
        slide.shapes.title.text = text


def add_body_text(slide, lines):
    """Add text to the body placeholder (index 1) if it exists."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.text = lines[0]
            for line in lines[1:]:
                p = tf.add_paragraph()
                p.text = line
            break


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Market Sizing Analysis"
    for ph in slide1.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = "FY2025 Strategic Planning | Prepared by Strategy Team"
            break

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide2, "Executive Summary")
    add_body_text(slide2, [
        "Total addressable market valued at $2B across North America and Europe",
        "Serviceable available market estimated at $800M based on product-market fit",
        "Target segment of $350M identified in enterprise SaaS vertical",
        "Current pipeline of $120M with 52% conversion rate expected",
        "Closed revenue of $62M represents 7.75% market penetration in target segment",
    ])

    # --- Slide 3: Market Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide3, "Market Overview")
    add_body_text(slide3, [
        "Cloud infrastructure spending projected to grow 22% YoY through 2026",
        "Enterprise adoption rate increased from 34% to 58% in past 18 months",
        "Key verticals: Financial Services, Healthcare, Manufacturing",
        "Competitive landscape: 3 major players control 65% of market share",
        "Regulatory tailwinds in EU and APAC driving compliance-related demand",
    ])

    # --- Slide 4: Competitive Landscape ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide4, "Competitive Landscape")
    add_body_text(slide4, [
        "AlphaTech Solutions - 28% market share, strong in financial services",
        "NovaSoft Inc. - 22% market share, dominant in healthcare vertical",
        "PeakCloud Systems - 15% market share, growing rapidly in manufacturing",
        "Our Position - 7.75% and accelerating with differentiated AI capabilities",
        "Remaining 27.25% fragmented among 40+ smaller vendors",
    ])

    # --- Slide 5: Revenue Projections ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide5, "Revenue Projections")
    add_body_text(slide5, [
        "Q1 2025: $14.2M (actual) - 8% above target",
        "Q2 2025: $16.8M (forecast) - driven by enterprise expansion deals",
        "Q3 2025: $15.5M (forecast) - seasonal adjustment factored in",
        "Q4 2025: $18.3M (forecast) - year-end budget flush expected",
        "Full Year 2025 Target: $64.8M representing 4.5% growth over FY2024",
    ])

    # --- Slide 6: Market Funnel (EMPTY - task target) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add only a title text box at the top
    txBox = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Market Funnel"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    # NO funnel shapes - this is the task for the agent

    # --- Slide 7: Go-to-Market Strategy ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide7, "Go-to-Market Strategy")
    add_body_text(slide7, [
        "Phase 1: Expand enterprise sales team by 40% in Q1-Q2",
        "Phase 2: Launch partner channel program targeting system integrators",
        "Phase 3: Geographic expansion into APAC markets starting Q3",
        "Investment: $8.5M in sales & marketing over next 12 months",
        "Expected ROI: 3.2x within 18 months based on pipeline conversion models",
    ])

    # --- Slide 8: Next Steps ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide8, "Next Steps & Timeline")
    add_body_text(slide8, [
        "Board approval for expansion budget by March 30, 2025",
        "Hire VP of APAC Operations by April 15, 2025",
        "Partner program launch event scheduled for May 2025",
        "Quarterly business review with updated pipeline metrics - June 2025",
        "Mid-year strategy adjustment based on H1 performance data",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
