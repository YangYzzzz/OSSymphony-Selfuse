"""
Reward Script: Funnel visualization on slide 6 with 5 rectangles
Task ID: impress_exec_062
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20) - 5 rectangle/auto shapes present on slide 6
  Component 2 (0.25) - Correct widths in decreasing order (10, 8, 6, 4, 2.5 in)
  Component 3 (0.20) - Correct fill colors for each rectangle
  Component 4 (0.20) - Correct text labels with white font
  Component 5 (0.15) - Horizontally centered and stacked vertically
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_062'

# Expected funnel data (top to bottom)
EXPECTED_FUNNEL = [
    {"label": "Total Addressable Market ($2B)", "width_in": 10.0, "color": "003366"},
    {"label": "Serviceable Market ($800M)",     "width_in": 8.0,  "color": "0D47A1"},
    {"label": "Target Segment ($350M)",         "width_in": 6.0,  "color": "1565C0"},
    {"label": "Pipeline ($120M)",               "width_in": 4.0,  "color": "1E88E5"},
    {"label": "Closed Revenue ($62M)",          "width_in": 2.5,  "color": "42A5F5"},
]

EMU_PER_INCH = 914400


def get_funnel_shapes(slide):
    """Extract rectangle/auto shapes from slide 6, excluding pre-existing placeholders and text boxes."""
    shapes = []
    for shape in slide.shapes:
        # AUTO_SHAPE (1) is the type for rectangles added programmatically
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            shapes.append(shape)
    # Sort by top position (vertical stacking order)
    shapes.sort(key=lambda s: s.top)
    return shapes


def verify_task(file_path):
    """
    Verify funnel visualization on slide 6.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[5]  # slide 6 (0-indexed)
    slide_width = prs.slide_width

    funnel_shapes = get_funnel_shapes(slide)
    num_funnel = len(funnel_shapes)

    # Component 1: 5 rectangle shapes present on slide 6 (0.20 points)
    try:
        if num_funnel == 5:
            print(f"PASS: Component 1 - Found 5 funnel shapes on slide 6 (0.20 pts)")
            total_score += 0.20
        elif num_funnel >= 3:
            partial = 0.10
            print(f"PARTIAL: Component 1 - Found {num_funnel}/5 funnel shapes (0.10 pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 - Found {num_funnel}/5 funnel shapes, need 5")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    if num_funnel == 0:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Correct widths in decreasing order (0.25 points)
    try:
        width_matches = 0
        widths_decreasing = True
        prev_width = float('inf')

        for i, shape in enumerate(funnel_shapes):
            width_in = shape.width / EMU_PER_INCH
            if i < len(EXPECTED_FUNNEL):
                expected_w = EXPECTED_FUNNEL[i]["width_in"]
                # 5% tolerance for width
                if abs(width_in - expected_w) / expected_w <= 0.05:
                    width_matches += 1
                    print(f"  Width {i}: {width_in:.2f}in matches expected {expected_w}in")
                else:
                    print(f"  Width {i}: {width_in:.2f}in does NOT match expected {expected_w}in")

            if width_in >= prev_width:
                widths_decreasing = False
            prev_width = width_in

        if width_matches == 5 and widths_decreasing:
            print(f"PASS: Component 2 - All 5 widths correct and decreasing (0.25 pts)")
            total_score += 0.25
        elif widths_decreasing and width_matches >= 3:
            partial = 0.15
            print(f"PARTIAL: Component 2 - {width_matches}/5 widths match, decreasing order OK ({partial} pts)")
            total_score += partial
        elif widths_decreasing:
            partial = 0.10
            print(f"PARTIAL: Component 2 - Widths are decreasing but only {width_matches}/5 match ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 - {width_matches}/5 widths match, decreasing={widths_decreasing}")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Correct fill colors (0.20 points)
    try:
        color_matches = 0
        for i, shape in enumerate(funnel_shapes):
            if i >= len(EXPECTED_FUNNEL):
                break
            expected_color = EXPECTED_FUNNEL[i]["color"].upper()
            try:
                fill = shape.fill
                if fill.type is not None and fill.type == 1:  # SOLID
                    actual_color = str(fill.fore_color.rgb).upper()
                    if actual_color == expected_color:
                        color_matches += 1
                        print(f"  Color {i}: {actual_color} matches expected {expected_color}")
                    else:
                        print(f"  Color {i}: {actual_color} does NOT match expected {expected_color}")
                else:
                    print(f"  Color {i}: Fill type is {fill.type}, expected SOLID (1)")
            except Exception as e:
                print(f"  Color {i}: Error reading fill - {e}")

        if color_matches == 5:
            print(f"PASS: Component 3 - All 5 fill colors correct (0.20 pts)")
            total_score += 0.20
        elif color_matches >= 3:
            partial = round(0.20 * color_matches / 5, 2)
            print(f"PARTIAL: Component 3 - {color_matches}/5 colors match ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 - Only {color_matches}/5 colors match")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Correct text labels with white font (0.20 points)
    try:
        label_matches = 0
        white_text_matches = 0

        for i, shape in enumerate(funnel_shapes):
            if i >= len(EXPECTED_FUNNEL):
                break
            expected_label = EXPECTED_FUNNEL[i]["label"]

            if shape.has_text_frame:
                actual_text = shape.text_frame.text.strip()
                # Check label text
                if actual_text == expected_label:
                    label_matches += 1
                    print(f"  Label {i}: '{actual_text}' matches")
                else:
                    print(f"  Label {i}: '{actual_text}' does NOT match '{expected_label}'")

                # Check white font color
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color.type is not None:
                                rgb = str(run.font.color.rgb).upper()
                                if rgb == "FFFFFF":
                                    white_text_matches += 1
                                    break
                        except Exception:
                            pass
                    else:
                        continue
                    break
            else:
                print(f"  Label {i}: Shape has no text frame")

        # Labels: 0.12 pts, white text: 0.08 pts
        label_score = 0.0
        if label_matches == 5:
            label_score = 0.12
            print(f"PASS: Component 4a - All 5 labels correct (0.12 pts)")
        elif label_matches >= 3:
            label_score = round(0.12 * label_matches / 5, 2)
            print(f"PARTIAL: Component 4a - {label_matches}/5 labels match ({label_score} pts)")
        else:
            print(f"FAIL: Component 4a - Only {label_matches}/5 labels match")

        white_score = 0.0
        if white_text_matches == 5:
            white_score = 0.08
            print(f"PASS: Component 4b - All 5 shapes have white text (0.08 pts)")
        elif white_text_matches >= 3:
            white_score = round(0.08 * white_text_matches / 5, 2)
            print(f"PARTIAL: Component 4b - {white_text_matches}/5 have white text ({white_score} pts)")
        else:
            print(f"FAIL: Component 4b - Only {white_text_matches}/5 have white text")

        total_score += label_score + white_score
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Horizontally centered and stacked vertically (0.15 points)
    try:
        centered_count = 0
        stacked_ok = True
        prev_top = -1

        slide_center = slide_width / 2

        for i, shape in enumerate(funnel_shapes):
            shape_center = shape.left + shape.width / 2
            # Check centering with 2% tolerance of slide width
            offset = abs(shape_center - slide_center)
            tolerance = slide_width * 0.02
            if offset <= tolerance:
                centered_count += 1
            else:
                print(f"  Center {i}: offset={offset/EMU_PER_INCH:.3f}in from slide center")

            # Check vertical stacking (each shape top > previous shape top)
            if shape.top <= prev_top:
                stacked_ok = False
                print(f"  Stack {i}: top={shape.top} not below prev top={prev_top}")
            prev_top = shape.top

        center_ok = centered_count == len(funnel_shapes)
        if center_ok and stacked_ok:
            print(f"PASS: Component 5 - All shapes centered and stacked vertically (0.15 pts)")
            total_score += 0.15
        elif stacked_ok and centered_count >= 3:
            partial = 0.10
            print(f"PARTIAL: Component 5 - {centered_count}/{len(funnel_shapes)} centered, stacked OK ({partial} pts)")
            total_score += partial
        elif stacked_ok:
            partial = 0.05
            print(f"PARTIAL: Component 5 - Stacked vertically but only {centered_count} centered ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 - centered={centered_count}/{len(funnel_shapes)}, stacked={stacked_ok}")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
