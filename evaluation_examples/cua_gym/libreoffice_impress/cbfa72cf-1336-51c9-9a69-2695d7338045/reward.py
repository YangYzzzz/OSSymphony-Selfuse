"""
Reward Script: Change rectangle border from solid to dashed (long-dash, 2pt, #7F8C8D)
Task ID: impress_ndo_057
Domain: libreoffice_impress
Scoring:
  Component 1: Dash pattern is lgDash (0.35 pts)
  Component 2: Line width is 2pt / 25400 EMU (0.25 pts)
  Component 3: Line color is #7F8C8D (0.25 pts)
  Component 4: Fill preserved at #ECF0F1 AND border changed (0.15 pts)
"""

import os
from lxml import etree

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_057'

# XML namespaces
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'


def find_rectangle(slide):
    """Find the rectangle shape (AUTO_SHAPE named 'Rectangle 2') on the slide."""
    for shape in slide.shapes:
        if shape.name == 'Rectangle 2' and str(shape.shape_type) == 'AUTO_SHAPE (1)':
            return shape
    # Fallback: find any rectangle auto-shape
    for shape in slide.shapes:
        prst = shape.element.find(f'.//{{{NS_A}}}prstGeom')
        if prst is not None and prst.get('prst') == 'rect' and str(shape.shape_type) == 'AUTO_SHAPE (1)':
            return shape
    return None


def get_line_properties(shape):
    """Extract line properties from shape XML."""
    ln = shape.element.find(f'.//{{{NS_A}}}ln')
    if ln is None:
        return None, None, None

    # Width
    width = ln.get('w')
    width_emu = int(width) if width else None

    # Dash style
    dash_elem = ln.find(f'{{{NS_A}}}prstDash')
    dash_val = dash_elem.get('val') if dash_elem is not None else None

    # Color
    color_elem = ln.find(f'.//{{{NS_A}}}srgbClr')
    color_val = color_elem.get('val') if color_elem is not None else None

    return width_emu, dash_val, color_val


def get_fill_color(shape):
    """Extract solid fill color from shape spPr."""
    spPr = shape.element.find(f'{{{NS_P}}}spPr')
    if spPr is None:
        spPr = shape.element.find(f'.//{{{NS_A}}}spPr')
    if spPr is None:
        return None
    fill = spPr.find(f'{{{NS_A}}}solidFill')
    if fill is None:
        return None
    clr = fill.find(f'{{{NS_A}}}srgbClr')
    return clr.get('val') if clr is not None else None


def persist_app_state():
    """Try to save any unsaved LibreOffice edits via Ctrl+S."""
    try:
        os.environ["DISPLAY"] = ":0"
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        import time
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Gate: must have at least 1 slide
    if len(prs.slides) < 1:
        print("FAIL: No slides found")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]
    rect = find_rectangle(slide)
    if rect is None:
        print("FAIL: No rectangle shape found on slide 1")
        print("REWARD: 0.0")
        return 0.0

    width_emu, dash_val, color_val = get_line_properties(rect)
    fill_color = get_fill_color(rect)

    print(f"DEBUG: line width={width_emu}, dash={dash_val}, color={color_val}, fill={fill_color}")

    # Component 1: Dash pattern is lgDash (0.35 points)
    # Initial: solid -> Golden: lgDash
    dash_ok = False
    try:
        if dash_val is not None and dash_val.lower() == 'lgdash':
            print(f"PASS: Component 1 - Dash pattern is lgDash (0.35 pts)")
            total_score += 0.35
            dash_ok = (dash_val.lower() == 'lgdash')
        else:
            print(f"FAIL: Component 1 - Expected dash=lgDash, found: {dash_val}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Line width is 2pt = 25400 EMU (0.25 points)
    # Initial: 12700 (1pt) -> Golden: 25400 (2pt)
    width_ok = False
    try:
        if width_emu is not None and abs(width_emu - 25400) <= 1000:
            print(f"PASS: Component 2 - Line width is ~2pt ({width_emu} EMU) (0.25 pts)")
            total_score += 0.25
            width_ok = (abs(width_emu - 25400) <= 1000)
        else:
            print(f"FAIL: Component 2 - Expected width ~25400 EMU (2pt), found: {width_emu}")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Line color is #7F8C8D (0.25 points)
    # Initial: 000000 (black) -> Golden: 7F8C8D
    color_ok = False
    try:
        if color_val is not None and color_val.upper() == '7F8C8D':
            print(f"PASS: Component 3 - Line color is #7F8C8D (0.25 pts)")
            total_score += 0.25
            color_ok = (color_val.upper() == '7F8C8D')
        else:
            print(f"FAIL: Component 3 - Expected line color 7F8C8D, found: {color_val}")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Fill preserved at #ECF0F1 AND at least one border property changed (0.15 points)
    # Compound check: fill must be unchanged AND task must have been attempted.
    # On initial_env, no border changes so this fails. On golden_env, both conditions met.
    any_border_changed = dash_ok or width_ok or color_ok
    try:
        fill_ok = fill_color is not None and fill_color.upper() == 'ECF0F1'
        if fill_ok and any_border_changed:
            print(f"PASS: Component 4 - Fill preserved at #ECF0F1 with border changes applied (0.15 pts)")
            total_score += 0.15
        elif not fill_ok:
            print(f"FAIL: Component 4 - Fill color changed or missing, found: {fill_color}")
        else:
            print(f"FAIL: Component 4 - No border changes detected, fill alone is a precondition")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
