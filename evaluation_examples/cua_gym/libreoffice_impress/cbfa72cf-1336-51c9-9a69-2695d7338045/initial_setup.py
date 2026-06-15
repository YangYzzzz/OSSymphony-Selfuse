"""
Initial Setup: Rectangle with solid border on slide 1
Task ID: impress_ndo_057
Domain: libreoffice_impress

Creates a presentation with one slide containing a rectangle (15cm x 10cm)
with a solid 1pt black border and #ECF0F1 fill.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_057'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Use a blank slide layout
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Add rectangle: 15cm x 10cm, centered on slide
    slide_w = prs.slide_width   # default 10 inches = 25.4 cm
    slide_h = prs.slide_height  # default 7.5 inches = 19.05 cm
    rect_w = Cm(15)
    rect_h = Cm(10)
    left = (slide_w - rect_w) // 2
    top = (slide_h - rect_h) // 2

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, rect_w, rect_h
    )

    # Fill: #ECF0F1 (light gray)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xEC, 0xF0, 0xF1)

    # Border: solid, 1pt, black
    line = shape.line
    line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    line.width = Pt(1)
    # Solid is the default dash style, but set explicitly
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    line.dash_style = MSO_LINE_DASH_STYLE.SOLID

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
