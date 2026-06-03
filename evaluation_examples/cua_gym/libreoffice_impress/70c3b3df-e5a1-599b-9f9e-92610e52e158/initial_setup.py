"""
Initial Setup: Team roster presentation with 12 slides, slide 9 has title only
Task ID: impress_rp_047
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_047'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_text(slide, text, left=Inches(0.5), top=Inches(0.3),
                   width=Inches(9), height=Inches(0.8),
                   font_size=Pt(32), bold=True, color=RGBColor(0x1A, 0x1A, 0x2E)):
    """Add a title textbox to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_body_text(slide, text, left=Inches(0.5), top=Inches(1.3),
                  width=Inches(9), height=Inches(4.5),
                  font_size=Pt(16), color=RGBColor(0x33, 0x33, 0x33)):
    """Add body text to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = font_size
    run.font.color.rgb = color
    return txBox


def add_bullet_list(slide, items, left=Inches(0.5), top=Inches(1.5),
                    width=Inches(9), height=Inches(5)):
    """Add a bulleted list to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(6)
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # Title Only or Blank

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_text(slide1, "Team Overview", font_size=Pt(40),
                   top=Inches(2.5), left=Inches(1), width=Inches(8))
    add_body_text(slide1, "Nexora Technologies Inc.\nQ2 2025 Organizational Review",
                  top=Inches(3.8), left=Inches(1), width=Inches(8),
                  font_size=Pt(20), color=RGBColor(0x55, 0x55, 0x55))

    # --- Slide 2: Company Vision ---
    slide2 = prs.slides.add_slide(blank_layout)
    add_title_text(slide2, "Our Vision")
    add_body_text(slide2, (
        "Nexora Technologies is dedicated to building intelligent automation "
        "solutions that transform how enterprises manage their digital workflows. "
        "Founded in 2019, we have grown from a 12-person startup to a 280-person "
        "organization spanning three continents."
    ))

    # --- Slide 3: Key Metrics ---
    slide3 = prs.slides.add_slide(blank_layout)
    add_title_text(slide3, "Key Performance Metrics")
    add_bullet_list(slide3, [
        "Annual Revenue: $47.2M (up 34% YoY)",
        "Customer Base: 1,240 enterprise accounts",
        "Employee Retention: 94.2%",
        "Net Promoter Score: 72",
        "Product Uptime: 99.97%",
        "Markets Served: 18 countries"
    ])

    # --- Slide 4: Department Overview ---
    slide4 = prs.slides.add_slide(blank_layout)
    add_title_text(slide4, "Department Structure")
    add_bullet_list(slide4, [
        "Engineering - 98 members across 8 squads",
        "Product & Design - 34 members, 3 product lines",
        "Sales & Marketing - 52 members, 4 regional teams",
        "Customer Success - 28 members, enterprise & SMB",
        "Operations & HR - 42 members",
        "Data Science & Analytics - 26 members"
    ])

    # --- Slide 5: Engineering Highlights ---
    slide5 = prs.slides.add_slide(blank_layout)
    add_title_text(slide5, "Engineering Highlights")
    add_bullet_list(slide5, [
        "Launched v3.0 platform with ML-driven automation engine",
        "Reduced average API response time from 280ms to 95ms",
        "Migrated 100% of infrastructure to Kubernetes",
        "Implemented zero-trust security architecture",
        "Open-sourced 3 internal libraries with 2,400+ GitHub stars"
    ])

    # --- Slide 6: Product Roadmap ---
    slide6 = prs.slides.add_slide(blank_layout)
    add_title_text(slide6, "Product Roadmap 2025")
    add_bullet_list(slide6, [
        "Q2: AI Assistant integration for workflow builder",
        "Q3: Real-time collaboration features",
        "Q3: Enterprise SSO & SCIM provisioning",
        "Q4: Mobile app launch (iOS & Android)",
        "Q4: Advanced analytics dashboard with custom reports"
    ])

    # --- Slide 7: Customer Success Stories ---
    slide7 = prs.slides.add_slide(blank_layout)
    add_title_text(slide7, "Customer Success Stories")
    add_bullet_list(slide7, [
        "Meridian Health Systems: 60% reduction in manual data entry",
        "Crestline Financial: Automated 85% of compliance reporting",
        "Vanguard Logistics: $2.1M annual savings from workflow automation",
        "Pinnacle Education: Onboarded 15,000 users in 3 weeks"
    ])

    # --- Slide 8: Culture & Values ---
    slide8 = prs.slides.add_slide(blank_layout)
    add_title_text(slide8, "Culture & Values")
    add_bullet_list(slide8, [
        "Innovation First: Dedicated 20% time for experimental projects",
        "Radical Transparency: All-hands meetings with open financials",
        "Customer Obsession: Every engineer does monthly support rotation",
        "Continuous Learning: $5,000 annual education budget per employee",
        "Global Mindset: Teams spanning San Francisco, London, and Singapore"
    ])

    # --- Slide 9: Core Team (BLANK except title) ---
    slide9 = prs.slides.add_slide(blank_layout)
    add_title_text(slide9, "Core Team", font_size=Pt(36))
    # Intentionally blank - the agent's task is to add team member cards here

    # --- Slide 10: Awards & Recognition ---
    slide10 = prs.slides.add_slide(blank_layout)
    add_title_text(slide10, "Awards & Recognition")
    add_bullet_list(slide10, [
        "2024 Gartner Cool Vendor in Intelligent Automation",
        "Inc. 5000 Fastest Growing Companies (#342)",
        "Best Places to Work - SF Business Times 2024",
        "G2 Leader in Workflow Automation (4 consecutive quarters)",
        "CIO 100 Award for Digital Transformation Excellence"
    ])

    # --- Slide 11: Financial Outlook ---
    slide11 = prs.slides.add_slide(blank_layout)
    add_title_text(slide11, "Financial Outlook")
    add_bullet_list(slide11, [
        "Projected FY2025 Revenue: $63M (target 34% growth)",
        "Series C Funding: $85M at $420M valuation",
        "Operating Margin Target: 18% by Q4 2025",
        "R&D Investment: 32% of revenue",
        "International Revenue: Growing from 22% to 35% of total"
    ])

    # --- Slide 12: Next Steps ---
    slide12 = prs.slides.add_slide(blank_layout)
    add_title_text(slide12, "Next Steps & Action Items")
    add_bullet_list(slide12, [
        "Complete Q2 hiring plan: 18 open engineering positions",
        "Finalize partnership with CloudBridge Solutions",
        "Launch internal mentorship program by June 15",
        "Submit SOC 2 Type II audit documentation",
        "Prepare board presentation for July quarterly review"
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
