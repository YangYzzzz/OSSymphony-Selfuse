"""
Reward Script: Team roster slide with 4 member cards
Task ID: impress_rp_047
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): 4 rounded rectangles with white fill and #DDDDDD border
  Component 2 (0.20): 4 ovals with alternating #3498DB / #9B59B6 colors
  Component 3 (0.20): 4 name text boxes (correct names, 16pt bold)
  Component 4 (0.20): 4 title text boxes (correct titles, 12pt #666666)
  Component 5 (0.20): 4 email text boxes (correct format, 10pt #5DADE2)
"""

import os

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_047'

# Expected card data
EXPECTED_NAMES = ['Sarah Chen', 'Marcus Williams', 'Priya Patel', "James O'Brien"]
EXPECTED_TITLES = ['Engineering Lead', 'Design Director', 'Product Manager', 'Data Scientist']
# Alternating circle colors: cards 1,3 = #3498DB; cards 2,4 = #9B59B6
EXPECTED_CIRCLE_COLORS = ['3498DB', '9B59B6', '3498DB', '9B59B6']


def get_font_color_rgb(run):
    """Safely get font color RGB string."""
    try:
        if run.font.color.type is not None:
            return str(run.font.color.rgb)
    except Exception:
        pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Check we have at least 9 slides
    if len(prs.slides) < 9:
        print(f"FAIL: Need at least 9 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[8]  # slide 9 (0-indexed)

    # Classify shapes on the slide
    rounded_rects = []
    ovals = []
    text_shapes = []

    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                ast = shape.auto_shape_type
                # ROUNDED_RECTANGLE = 5, OVAL = 9
                if ast == 5:
                    rounded_rects.append(shape)
                elif ast == 9:
                    ovals.append(shape)
        except Exception:
            pass
        if hasattr(shape, 'text_frame') and shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            text = shape.text_frame.text.strip()
            if text:
                text_shapes.append(shape)

    # Also check group shapes recursively
    def extract_from_groups(shape):
        rr, ov, ts = [], [], []
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                try:
                    if sub.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        ast = sub.auto_shape_type
                        if ast == 5:
                            rr.append(sub)
                        elif ast == 9:
                            ov.append(sub)
                except Exception:
                    pass
                if hasattr(sub, 'text_frame') and hasattr(sub, 'text'):
                    t = sub.text_frame.text.strip()
                    if t:
                        ts.append(sub)
                r2, o2, t2 = extract_from_groups(sub)
                rr.extend(r2)
                ov.extend(o2)
                ts.extend(t2)
        return rr, ov, ts

    for shape in slide.shapes:
        if hasattr(shape, 'shapes'):
            rr, ov, ts = extract_from_groups(shape)
            rounded_rects.extend(rr)
            ovals.extend(ov)
            text_shapes.extend(ts)

    print(f"Found: {len(rounded_rects)} rounded rects, {len(ovals)} ovals, {len(text_shapes)} text shapes")

    # Component 1: 4 rounded rectangles with white fill and #DDDDDD border (0.20 points)
    try:
        valid_rects = 0
        for rect in rounded_rects:
            try:
                # Check white fill
                has_white_fill = (rect.fill.type == 1 and str(rect.fill.fore_color.rgb) == 'FFFFFF')
                # Check DDDDDD border
                has_dd_border = False
                try:
                    has_dd_border = (rect.line.fill.type == 1 and str(rect.line.color.rgb) == 'DDDDDD')
                except Exception:
                    pass
                if has_white_fill and has_dd_border:
                    valid_rects += 1
            except Exception:
                pass
        if valid_rects >= 4:
            print(f"PASS: Component 1 — {valid_rects} valid rounded rectangles with white fill and #DDDDDD border (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 — Expected 4 valid rounded rects, found {valid_rects}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: 4 ovals with correct alternating colors (0.20 points)
    try:
        # Sort ovals by left position to determine order
        sorted_ovals = sorted(ovals, key=lambda s: s.left)
        oval_colors = []
        for oval in sorted_ovals:
            try:
                if oval.fill.type == 1:  # SOLID
                    oval_colors.append(str(oval.fill.fore_color.rgb))
                else:
                    oval_colors.append(None)
            except Exception:
                oval_colors.append(None)

        if len(sorted_ovals) >= 4:
            # Check alternating pattern: 3498DB, 9B59B6, 3498DB, 9B59B6
            color_match = 0
            for idx in range(4):
                if oval_colors[idx] == EXPECTED_CIRCLE_COLORS[idx]:
                    color_match += 1
            if color_match == 4:
                print(f"PASS: Component 2 — 4 ovals with correct alternating colors {oval_colors[:4]} (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 2 — Oval colors {oval_colors[:4]} vs expected {EXPECTED_CIRCLE_COLORS}, matched {color_match}/4")
        else:
            print(f"FAIL: Component 2 — Expected 4 ovals, found {len(sorted_ovals)}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Filter text shapes: exclude the slide title "Core Team"
    card_text_shapes = [s for s in text_shapes if s.text_frame.text.strip() != 'Core Team']

    # Categorize text shapes by their content
    name_shapes = []
    title_shapes = []
    email_shapes = []

    for ts in card_text_shapes:
        text = ts.text_frame.text.strip()
        # Check if it matches a known name
        if text in EXPECTED_NAMES:
            name_shapes.append(ts)
        elif text in EXPECTED_TITLES:
            title_shapes.append(ts)
        elif '@' in text:
            email_shapes.append(ts)

    # Component 3: 4 name text boxes with correct names, 16pt bold (0.20 points)
    try:
        valid_names = 0
        found_names = set()
        for ns in name_shapes:
            text = ns.text_frame.text.strip()
            runs = [r for r in ns.text_frame.paragraphs[0].runs if (r.text or "").strip()]
            if runs:
                run = runs[0]
                size_ok = run.font.size is not None and abs(run.font.size - 203200) < 5000  # ~16pt tolerance
                bold_ok = run.font.bold is True
                if size_ok and bold_ok and text not in found_names:
                    valid_names += 1
                    found_names.add(text)
                    print(f"  Name OK: {text} (size={run.font.size}, bold={run.font.bold})")
                else:
                    print(f"  Name partial: {text} (size={run.font.size}, bold={run.font.bold}, size_ok={size_ok}, bold_ok={bold_ok})")
        if valid_names >= 4:
            print(f"PASS: Component 3 — {valid_names} names with correct 16pt bold formatting (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 3 — Expected 4 valid names, found {valid_names}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: 4 title text boxes with correct titles, 12pt #666666 (0.20 points)
    try:
        valid_titles = 0
        found_titles = set()
        for ts in title_shapes:
            text = ts.text_frame.text.strip()
            runs = [r for r in ts.text_frame.paragraphs[0].runs if (r.text or "").strip()]
            if runs:
                run = runs[0]
                size_ok = run.font.size is not None and abs(run.font.size - 152400) < 5000  # ~12pt tolerance
                color = get_font_color_rgb(run)
                color_ok = color == '666666'
                if size_ok and color_ok and text not in found_titles:
                    valid_titles += 1
                    found_titles.add(text)
                    print(f"  Title OK: {text} (size={run.font.size}, color={color})")
                else:
                    print(f"  Title partial: {text} (size={run.font.size}, color={color}, size_ok={size_ok}, color_ok={color_ok})")
        if valid_titles >= 4:
            print(f"PASS: Component 4 — {valid_titles} titles with correct 12pt #666666 formatting (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 4 — Expected 4 valid titles, found {valid_titles}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: 4 email text boxes with correct format, 10pt #5DADE2 (0.20 points)
    try:
        valid_emails = 0
        found_emails = set()
        for es in email_shapes:
            text = es.text_frame.text.strip()
            runs = [r for r in es.text_frame.paragraphs[0].runs if (r.text or "").strip()]
            if runs:
                run = runs[0]
                size_ok = run.font.size is not None and abs(run.font.size - 127000) < 5000  # ~10pt tolerance
                color = get_font_color_rgb(run)
                color_ok = color == '5DADE2'
                if size_ok and color_ok and text not in found_emails:
                    valid_emails += 1
                    found_emails.add(text)
                    print(f"  Email OK: {text} (size={run.font.size}, color={color})")
                else:
                    print(f"  Email partial: {text} (size={run.font.size}, color={color}, size_ok={size_ok}, color_ok={color_ok})")
        if valid_emails >= 4:
            print(f"PASS: Component 5 — {valid_emails} emails with correct 10pt #5DADE2 formatting (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 5 — Expected 4 valid emails, found {valid_emails}")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
