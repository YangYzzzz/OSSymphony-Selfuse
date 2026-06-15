"""
Reward Script: Insert remarks from remarks.docx into notes fields of Annual_Report.pptx
Task ID: osworld_multi_apps_impress_notes_import_002
Domain: libreoffice_impress
Scoring:
  - Component 1: All 8 slides have non-empty notes (0.2 points)
  - Component 2: Each slide's notes text matches the expected remark from remarks.docx (0.8 points, 0.1 per slide)
Total: 1.0
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_multi_apps_impress_notes_import_002'

# Expected notes for each slide (1-indexed), derived from remarks.docx
EXPECTED_NOTES = {
    1: "Welcome everyone to Meridian Technologies' Annual Report presentation for fiscal year 2024. This year has been transformational for our company, and I'm excited to walk you through our key achievements and future plans. Please hold questions until the end.",
    2: "As you can see from our executive summary, 2024 was a record-breaking year. The 18% revenue growth exceeded our initial guidance of 14-16%. Our international expansion strategy is clearly paying dividends, and our workforce growth reflects sustained confidence in the business.",
    3: "Let me highlight a few key financial metrics. Our gross margin of 64.4% is the highest in company history, driven by the shift toward higher-margin SaaS products. EBITDA of $33.2M gives us significant capacity for reinvestment and potential M&A activity. The EPS growth of 18.6% reflects strong shareholder value creation.",
    4: "CloudSync Pro continues to be our star performer with 32% growth. The enterprise segment is showing strong adoption, particularly among financial services clients. We've made a strategic decision to wind down legacy products by end of 2025, which will free resources for higher-growth areas.",
    5: "North America remains our core market, but international growth is accelerating. Europe grew 28% this year, outpacing our overall growth rate. The new regions we entered in Q3 are already generating meaningful revenue, ahead of our 18-month break-even projection.",
    6: "Our R&D investment at 13% of revenue is in line with industry leaders. The AI analytics module has received very positive feedback from beta customers, and we expect to launch it generally in Q1 2025. The Stanford partnership will accelerate our machine learning capabilities significantly.",
    7: "Our 2025 guidance reflects continued strong momentum. The $165M-$172M range assumes normal macro conditions and successful CloudSync Pro v3.0 launch. We're being deliberately conservative given global economic uncertainties. The 25%+ operating margin target is achievable through continued operating leverage.",
    8: "Thank you all for your time and continued support of Meridian Technologies. We are proud of what we've accomplished in 2024 and remain focused on executing our strategy. I'll now open the floor for questions. Please keep questions concise so we can address as many as possible.",
}


def get_slide_notes(slide):
    """Retrieve notes text from a slide, stripping whitespace."""
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: verify we have 8 slides
    num_slides = len(prs.slides)
    if num_slides != 8:
        print(f"CRITICAL: Expected 8 slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Presentation loaded — {num_slides} slides")

    # Component 1: All 8 slides have non-empty notes (0.2 points)
    # This FAILS on initial_env (all empty) and PASSES on golden_env (all populated)
    try:
        empty_slides = [
            i + 1 for i, slide in enumerate(prs.slides)
            if not get_slide_notes(slide)
        ]

        if len(empty_slides) == 0:
            print(f"PASS: Component 1 — All 8 slides have non-empty notes (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 1 — Slides {empty_slides} have empty notes")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Each slide's notes match the expected remarks text (0.8 points total, 0.1 per slide)
    # The expected text is extracted from remarks.docx (Slide N: labels map to slides)
    points_per_slide = 0.1
    for slide_num in range(1, 9):
        try:
            slide = prs.slides[slide_num - 1]
            actual_notes = get_slide_notes(slide)
            expected_notes = EXPECTED_NOTES[slide_num]

            # Normalize whitespace for comparison
            actual_normalized = " ".join(actual_notes.split())
            expected_normalized = " ".join(expected_notes.split())

            if actual_normalized == expected_normalized:
                print(f"PASS: Component 2 Slide {slide_num} — notes match expected text ({points_per_slide} pts)")
                total_score += points_per_slide
            else:
                # Partial credit: check if the core content is present (starts with expected opening)
                # We require exact content match for full score
                print(f"FAIL: Component 2 Slide {slide_num} — notes do not match")
                print(f"  Expected: \"{expected_normalized[:80]}...\"")
                print(f"  Actual:   \"{actual_normalized[:80]}...\"")
        except Exception as e:
            print(f"ERROR: Component 2 Slide {slide_num} — {e}")

    final_score = round(min(total_score, 1.0), 4)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path on the VM
file_path = f'{WORKDIR}/Annual_Report.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
