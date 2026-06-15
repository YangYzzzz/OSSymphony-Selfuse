"""
Initial Setup: Insert remarks from Word document into Impress slide notes
Task ID: osworld_multi_apps_impress_notes_import_002
Domain: libreoffice_impress (multi-app: also creates remarks.docx)

Creates:
  - /home/user/Annual_Report.pptx  (8 slides, NO notes on any slide)
  - /home/user/Desktop/remarks.docx (remarks labeled 'Slide N:' for slides 1-8)

Then opens Annual_Report.pptx in LibreOffice Impress.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from docx import Document
from docx.shared import Pt as DocxPt

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_impress_notes_import_002'
PPTX_OUTPUT = f'{WORKDIR}/Annual_Report.pptx'
DOCX_OUTPUT = f'{DESKTOP}/remarks.docx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


# Realistic slide content for an Annual Report
SLIDES = [
    {
        "title": "Annual Report 2024",
        "subtitle": "Meridian Technologies Inc.\nFiscal Year Overview",
        "layout": 0,  # Title Slide
    },
    {
        "title": "Executive Summary",
        "body": (
            "• Revenue grew 18% year-over-year to $142M\n"
            "• Operating margins improved to 23.4%\n"
            "• Expanded into 3 new international markets\n"
            "• Headcount increased from 580 to 724 employees"
        ),
        "layout": 1,
    },
    {
        "title": "Financial Highlights",
        "body": (
            "• Total Revenue: $142.3 Million\n"
            "• Gross Profit: $91.7 Million (64.4% margin)\n"
            "• EBITDA: $33.2 Million\n"
            "• Net Income: $21.8 Million\n"
            "• EPS: $2.74 (up from $2.31)"
        ),
        "layout": 1,
    },
    {
        "title": "Product Performance",
        "body": (
            "• CloudSync Pro: $58M revenue (+32%)\n"
            "• DataBridge Enterprise: $47M revenue (+11%)\n"
            "• SecureVault Suite: $24M revenue (+8%)\n"
            "• Legacy Products: $13M revenue (-14%)\n"
            "• New product launches: Q2 and Q4"
        ),
        "layout": 1,
    },
    {
        "title": "Regional Performance",
        "body": (
            "• North America: $89M (62.6% of total)\n"
            "• Europe: $31M (21.8% of total)\n"
            "• Asia-Pacific: $16M (11.2% of total)\n"
            "• Rest of World: $6.3M (4.4% of total)\n"
            "• New regions: Brazil, South Korea, Netherlands"
        ),
        "layout": 1,
    },
    {
        "title": "R&D and Innovation",
        "body": (
            "• R&D Investment: $18.5M (13% of revenue)\n"
            "• 47 new patents filed in 2024\n"
            "• AI-powered analytics module in beta\n"
            "• Partnership with Stanford AI Lab announced\n"
            "• 3 product roadmap milestones exceeded"
        ),
        "layout": 1,
    },
    {
        "title": "Outlook for 2025",
        "body": (
            "• Revenue guidance: $165M–$172M\n"
            "• Target operating margin: 25%+\n"
            "• Planned headcount additions: ~150\n"
            "• Expansion into Southeast Asia and Middle East\n"
            "• CloudSync Pro v3.0 launch scheduled for Q2"
        ),
        "layout": 1,
    },
    {
        "title": "Thank You",
        "subtitle": "Questions & Discussion\n\ncontact@meridiantech.com",
        "layout": 0,  # Title Slide (closing)
    },
]

# Corresponding remarks for each slide (will go into remarks.docx)
REMARKS = [
    (
        "Welcome everyone to Meridian Technologies' Annual Report presentation for fiscal year 2024. "
        "This year has been transformational for our company, and I'm excited to walk you through "
        "our key achievements and future plans. Please hold questions until the end."
    ),
    (
        "As you can see from our executive summary, 2024 was a record-breaking year. "
        "The 18% revenue growth exceeded our initial guidance of 14-16%. "
        "Our international expansion strategy is clearly paying dividends, "
        "and our workforce growth reflects sustained confidence in the business."
    ),
    (
        "Let me highlight a few key financial metrics. Our gross margin of 64.4% is the highest "
        "in company history, driven by the shift toward higher-margin SaaS products. "
        "EBITDA of $33.2M gives us significant capacity for reinvestment and potential M&A activity. "
        "The EPS growth of 18.6% reflects strong shareholder value creation."
    ),
    (
        "CloudSync Pro continues to be our star performer with 32% growth. "
        "The enterprise segment is showing strong adoption, particularly among financial services clients. "
        "We've made a strategic decision to wind down legacy products by end of 2025, "
        "which will free resources for higher-growth areas."
    ),
    (
        "North America remains our core market, but international growth is accelerating. "
        "Europe grew 28% this year, outpacing our overall growth rate. "
        "The new regions we entered in Q3 are already generating meaningful revenue, "
        "ahead of our 18-month break-even projection."
    ),
    (
        "Our R&D investment at 13% of revenue is in line with industry leaders. "
        "The AI analytics module has received very positive feedback from beta customers, "
        "and we expect to launch it generally in Q1 2025. "
        "The Stanford partnership will accelerate our machine learning capabilities significantly."
    ),
    (
        "Our 2025 guidance reflects continued strong momentum. "
        "The $165M-$172M range assumes normal macro conditions and successful CloudSync Pro v3.0 launch. "
        "We're being deliberately conservative given global economic uncertainties. "
        "The 25%+ operating margin target is achievable through continued operating leverage."
    ),
    (
        "Thank you all for your time and continued support of Meridian Technologies. "
        "We are proud of what we've accomplished in 2024 and remain focused on executing our strategy. "
        "I'll now open the floor for questions. Please keep questions concise so we can address as many as possible."
    ),
]


def create_initial_pptx():
    """Create Annual_Report.pptx with 8 slides and NO notes."""
    prs = Presentation()

    # Set standard widescreen dimensions
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for i, slide_data in enumerate(SLIDES):
        layout_idx = slide_data.get("layout", 1)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data["title"]

        # Set body/subtitle
        if layout_idx == 0:
            # Title slide: placeholders[1] is subtitle
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_data.get("subtitle", "")
        else:
            # Content slide: placeholders[1] is body
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_data.get("body", "")

        # IMPORTANT: Do NOT set notes — initial state must have no notes

    prs.save(PPTX_OUTPUT)
    print(f'Initial PPTX created (no notes): {PPTX_OUTPUT}')


def create_remarks_docx():
    """Create remarks.docx on Desktop with 'Slide N:' labeled remarks."""
    os.makedirs(DESKTOP, exist_ok=True)

    doc = Document()

    # Add a title
    title_para = doc.add_heading('Presenter Remarks — Annual Report 2024', level=1)

    doc.add_paragraph('')  # spacer

    for i, remark_text in enumerate(REMARKS, 1):
        # Add slide label as bold paragraph
        label_para = doc.add_paragraph()
        label_run = label_para.add_run(f'Slide {i}:')
        label_run.bold = True
        label_run.font.size = DocxPt(12)

        # Add remark text
        remark_para = doc.add_paragraph(remark_text)
        remark_para.paragraph_format.left_indent = DocxPt(12)

        # Spacer between slides
        doc.add_paragraph('')

    doc.save(DOCX_OUTPUT)
    print(f'Remarks DOCX created: {DOCX_OUTPUT}')


def create_initial():
    create_initial_pptx()
    create_remarks_docx()

    # GUI-ready startup: open Annual_Report.pptx in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{PPTX_OUTPUT}"', delay_sec=3.0)
    print('GUI_READY: launched LibreOffice Impress with Annual_Report.pptx (DISPLAY=:0)')


create_initial()
