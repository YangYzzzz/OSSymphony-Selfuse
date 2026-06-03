"""
Initial Setup: Create a 7-slide presentation template for deletion/insertion task
Task ID: impstruct_010
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impstruct_010'
OUTPUT = f'{WORKDIR}/old_template.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Cover ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only
    txBox = slide1.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Cover"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    # Subtitle
    txBox2 = slide1.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Meridian Solutions Inc. — Annual Strategy Review 2025"
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.name = "Arial"
    run2.font.size = Pt(20)
    run2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # Background
    fill1 = slide1.background.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Executive Summary"
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    body2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11), Inches(5))
    tf2 = body2.text_frame
    tf2.word_wrap = True
    items = [
        "Revenue grew 18% year-over-year to $42.3M in FY2024",
        "Customer retention rate improved from 87% to 93%",
        "Launched 3 new product lines across APAC markets",
        "Headcount increased by 45 employees (Engineering, Sales)",
        "Operating margin expanded from 12.4% to 15.1%",
    ]
    for i, item in enumerate(items):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f"• {item}"
        p.space_after = Pt(10)
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(18)

    # --- Slide 3: Financial Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Financial Overview — Q4 2024"
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    # Table
    table_shape = slide3.shapes.add_table(5, 4, Inches(1), Inches(1.8), Inches(10), Inches(3))
    table = table_shape.table
    headers = ["Metric", "Q3 2024", "Q4 2024", "Change"]
    data = [
        ["Revenue", "$9.8M", "$11.2M", "+14.3%"],
        ["COGS", "$5.1M", "$5.6M", "+9.8%"],
        ["Gross Profit", "$4.7M", "$5.6M", "+19.1%"],
        ["Net Income", "$1.2M", "$1.7M", "+41.7%"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 4: Regional Performance ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide4.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Regional Performance Breakdown"
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    regions = [
        ("North America", "$22.1M", "52% of total revenue, +12% YoY"),
        ("Europe", "$11.8M", "28% of total revenue, +21% YoY"),
        ("Asia-Pacific", "$6.3M", "15% of total revenue, +35% YoY"),
        ("Rest of World", "$2.1M", "5% of total revenue, +8% YoY"),
    ]
    body4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11), Inches(5))
    tf4 = body4.text_frame
    tf4.word_wrap = True
    for i, (region, rev, detail) in enumerate(regions):
        if i == 0:
            p = tf4.paragraphs[0]
        else:
            p = tf4.add_paragraph()
        p.text = f"{region}: {rev}"
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(20)
        run.font.bold = True
        p2 = tf4.add_paragraph()
        p2.text = f"   {detail}"
        run2 = p2.runs[0]
        run2.font.name = "Arial"
        run2.font.size = Pt(16)
        run2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # --- Slide 5: Product Roadmap ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide5.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Product Roadmap — 2025"
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    milestones = [
        "Q1: Launch DataSync Pro v3.0 with real-time collaboration",
        "Q2: Mobile app redesign (iOS + Android)",
        "Q3: Enterprise SSO and compliance dashboard",
        "Q4: AI-powered analytics module beta release",
    ]
    body5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11), Inches(5))
    tf5 = body5.text_frame
    tf5.word_wrap = True
    for i, item in enumerate(milestones):
        if i == 0:
            p = tf5.paragraphs[0]
        else:
            p = tf5.add_paragraph()
        p.text = f"▸ {item}"
        p.space_after = Pt(12)
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(18)

    # --- Slide 6: Team & Hiring ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide6.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Team Growth & Hiring Plan"
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    table_shape6 = slide6.shapes.add_table(5, 3, Inches(1.5), Inches(1.8), Inches(9), Inches(3))
    table6 = table_shape6.table
    headers6 = ["Department", "Current Headcount", "2025 Target"]
    data6 = [
        ["Engineering", "78", "105"],
        ["Sales & Marketing", "34", "48"],
        ["Operations", "22", "28"],
        ["Customer Success", "16", "24"],
    ]
    for c, h in enumerate(headers6):
        cell = table6.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row_data in enumerate(data6, 1):
        for c, val in enumerate(row_data):
            table6.cell(r, c).text = val

    # --- Slide 7: Back Cover ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide7.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Back Cover"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    txBox2 = slide7.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Thank you for your attention."
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.name = "Arial"
    run2.font.size = Pt(20)
    run2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p3 = tf2.add_paragraph()
    p3.text = "Contact: strategy@meridiansolutions.com"
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.runs[0]
    run3.font.name = "Arial"
    run3.font.size = Pt(16)
    run3.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    fill7 = slide7.background.fill
    fill7.solid()
    fill7.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
