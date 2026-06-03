"""
Reward Script: Delete all slides except the first and last, then insert 3 new blank slides between them.
Task ID: impstruct_010
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): Slide count changed from 7 to exactly 5
  Component 2 (0.25): First slide is 'Cover' AND last slide is 'Back Cover' AND count is 5
                       (verifies correct slides were kept after deletion)
  Component 3 (0.25): Middle slides (2-4) are blank (no text content) AND count is 5
  Component 4 (0.2): Old content slides (Executive Summary, Financial, Regional, Roadmap, Team)
                      are no longer present in the presentation
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impstruct_010'


def persist_app_state(domain):
    """Send Ctrl+S to save any unsaved GUI edits."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print(f"PERSIST: ctrl+s sent for {domain}")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def get_slide_text(slide):
    """Extract all non-empty text from a slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return texts


def get_all_text(prs):
    """Get all text across all slides."""
    all_text = []
    for slide in prs.slides:
        all_text.extend(get_slide_text(slide))
    return all_text


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Slide count is exactly 5 (0.3 points)
    # Initial has 7 slides; after task it should be 5.
    try:
        if num_slides == 5:
            print(f"PASS: Component 1 - Slide count is 5 (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 - Expected 5 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: First slide is 'Cover' AND last slide is 'Back Cover' AND count is 5 (0.25 pts)
    # This compound check ensures the correct slides were preserved after deletion.
    # On initial (7 slides), this fails because count != 5.
    try:
        if num_slides == 5:
            first_texts = get_slide_text(prs.slides[0])
            last_texts = get_slide_text(prs.slides[-1])
            has_cover = any("Cover" in t and "Back Cover" not in t for t in first_texts)
            has_back_cover = any("Back Cover" in t for t in last_texts)
            if has_cover and has_back_cover:
                print(f"PASS: Component 2 - Cover/Back Cover preserved with 5 slides (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 2 - First texts: {first_texts}, Last texts: {last_texts}")
        else:
            print(f"FAIL: Component 2 - Slide count is {num_slides}, not 5")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Middle slides (2-4) are blank - no text content (0.25 points)
    # Only valid when there are exactly 5 slides.
    try:
        if num_slides == 5:
            non_blank = [idx for idx in [1, 2, 3]
                         if get_slide_text(prs.slides[idx])]
            for idx in non_blank:
                print(f"  Slide {idx+1} has text: {get_slide_text(prs.slides[idx])}")
            if not non_blank:
                print(f"PASS: Component 3 - Slides 2-4 are blank (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 3 - Some middle slides have text content")
        else:
            print(f"FAIL: Component 3 - Cannot verify middle slides (count is {num_slides})")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Old content slides are gone (0.2 points)
    # The deleted slides had unique text: "Executive Summary", "Financial Overview",
    # "Regional Performance", "Product Roadmap", "Team Growth".
    # None of these should appear anywhere in the final presentation.
    try:
        all_text = " ".join(get_all_text(prs))
        deleted_markers = [
            "Executive Summary",
            "Financial Overview",
            "Regional Performance",
            "Product Roadmap",
            "Team Growth",
        ]
        found_old = [m for m in deleted_markers if m in all_text]
        if not found_old:
            print(f"PASS: Component 4 - Old content slides removed (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 4 - Old content still present: {found_old}")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/old_template.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
