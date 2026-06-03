"""
Reward Script: Jeopardy-style game board on slide 2
Task ID: impress_teach_060
Domain: libreoffice_impress
Scoring:
  - Component 1: Table exists on slide 2 with correct dimensions (0.15)
  - Component 2: Header row has correct category names (0.20)
  - Component 3: Data rows have correct point values (0.20)
  - Component 4: Header row background color #1A237E with white bold text (0.25)
  - Component 5: Data cells background color #283593 with white text (0.20)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_060'

def persist_app_state(domain):
    """Save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(1.0)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: need at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 2")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]

    # Find the table on slide 2
    table_shape = None
    for shape in slide2.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_shape = shape
            break

    if table_shape is None:
        print("FAIL: No table found on slide 2")
        print("REWARD: 0.0")
        return 0.0

    table = table_shape.table
    num_rows = len(table.rows)
    num_cols = len(table.columns)

    # Component 1: Table dimensions (0.15 points)
    # Task says "6x5 table" but content describes 5 categories + 4 point rows = 5 rows.
    # Accept 5x5 (matching content) or 6x5 (matching literal dimension).
    try:
        if num_cols == 5 and num_rows >= 5:
            print(f"PASS: Component 1 — Table is {num_rows}x{num_cols} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — Table is {num_rows}x{num_cols}, expected 5+ rows x 5 cols")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Header row has correct category names (0.20 points)
    expected_headers = ['Vocabulary', 'Dates', 'People', 'Events', 'Geography']
    try:
        if num_cols >= 5:
            actual_headers = [table.cell(0, c).text.strip() for c in range(5)]
            matching = sum(1 for a, e in zip(actual_headers, expected_headers) if a.lower() == e.lower())
            if matching == 5:
                print(f"PASS: Component 2 — All 5 category headers correct: {actual_headers} (0.20 pts)")
                total_score += 0.20
            elif matching >= 3:
                partial = round(0.20 * (matching / 5), 2)
                print(f"PARTIAL: Component 2 — {matching}/5 headers correct: {actual_headers} ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 2 — Headers: {actual_headers}, expected: {expected_headers}")
        else:
            print(f"FAIL: Component 2 — Not enough columns ({num_cols})")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Data rows have correct point values (0.20 points)
    expected_values = ['100', '200', '300', '400']
    try:
        data_rows = num_rows - 1  # exclude header
        rows_correct = 0
        for r_idx, expected_val in enumerate(expected_values):
            row = r_idx + 1
            if row >= num_rows:
                break
            row_texts = [table.cell(row, c).text.strip() for c in range(min(num_cols, 5))]
            if all(t == expected_val for t in row_texts):
                rows_correct += 1
            else:
                print(f"  INFO: Row {row} values: {row_texts}, expected all '{expected_val}'")

        if rows_correct == 4:
            print(f"PASS: Component 3 — All 4 data rows have correct point values (0.20 pts)")
            total_score += 0.20
        elif rows_correct >= 2:
            partial = round(0.20 * (rows_correct / 4), 2)
            print(f"PARTIAL: Component 3 — {rows_correct}/4 data rows correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — Only {rows_correct}/4 data rows correct")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Header row formatting — #1A237E bg, white bold text (0.25 points)
    try:
        header_format_ok = 0
        for c in range(min(num_cols, 5)):
            cell = table.cell(0, c)
            # Check fill color
            fill_ok = False
            try:
                if cell.fill.type is not None:
                    fill_rgb = str(cell.fill.fore_color.rgb).upper()
                    fill_ok = (fill_rgb == '1A237E')  # derived from API check
            except:
                pass

            # Check text: white (#FFFFFF) and bold
            text_ok = False
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    if run.text.strip():
                        is_bold = run.font.bold is True
                        try:
                            font_rgb = str(run.font.color.rgb).upper()
                            is_white = font_rgb == 'FFFFFF'
                        except:
                            is_white = False
                        if is_bold and is_white:
                            text_ok = not False  # derived from conditional check

            if fill_ok and text_ok:
                header_format_ok += 1

        if header_format_ok == 5:
            print(f"PASS: Component 4 — All 5 header cells: #1A237E bg, white bold text (0.25 pts)")
            total_score += 0.25
        elif header_format_ok >= 3:
            partial = round(0.25 * (header_format_ok / 5), 2)
            print(f"PARTIAL: Component 4 — {header_format_ok}/5 header cells formatted correctly ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — Only {header_format_ok}/5 header cells formatted correctly")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Data cell formatting — #283593 bg, white text (0.20 points)
    try:
        total_data_cells = 0
        data_format_ok = 0
        for r in range(1, min(num_rows, 5)):  # rows 1-4
            for c in range(min(num_cols, 5)):
                total_data_cells += 1
                cell = table.cell(r, c)
                # Check fill color
                fill_ok = False
                try:
                    if cell.fill.type is not None:
                        fill_rgb = str(cell.fill.fore_color.rgb).upper()
                        fill_ok = (fill_rgb == '283593')  # derived from API check
                except:
                    pass

                # Check text: white (#FFFFFF)
                text_ok = False
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        if run.text.strip():
                            try:
                                font_rgb = str(run.font.color.rgb).upper()
                                text_ok = (font_rgb == 'FFFFFF') or text_ok  # derived from API check
                            except:
                                pass

                if fill_ok and text_ok:
                    data_format_ok += 1

        if total_data_cells > 0 and data_format_ok == total_data_cells:
            print(f"PASS: Component 5 — All {total_data_cells} data cells: #283593 bg, white text (0.20 pts)")
            total_score += 0.20
        elif total_data_cells > 0 and data_format_ok >= total_data_cells * 0.6:
            partial = round(0.20 * (data_format_ok / total_data_cells), 2)
            print(f"PARTIAL: Component 5 — {data_format_ok}/{total_data_cells} data cells formatted correctly ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 — Only {data_format_ok}/{total_data_cells} data cells formatted correctly")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
