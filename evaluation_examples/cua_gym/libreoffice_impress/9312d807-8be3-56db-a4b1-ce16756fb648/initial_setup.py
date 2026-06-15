"""
Initial Setup: Create a 3-slide Review_Game presentation. Slide 2 has title 'History Jeopardy' but no table.
Task ID: impress_teach_060
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_060'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "History Review Game"
    slide1.placeholders[1].text = "World History - Semester Final Review"

    # --- Slide 2: Title Only - 'History Jeopardy' (NO table, agent must add it) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "History Jeopardy"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    # --- Slide 3: Instructions slide ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Game Rules"
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.runs[0]
    run3.font.size = Pt(32)
    run3.font.bold = True

    rules_box = slide3.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    rtf = rules_box.text_frame
    rtf.word_wrap = True
    rules = [
        "1. Teams take turns selecting a category and point value.",
        "2. The teacher reads the answer; teams must respond with the correct question.",
        "3. Correct responses earn the point value; incorrect responses lose points.",
        "4. The team with the most points at the end wins.",
        "5. Daily Doubles are hidden randomly and allow teams to wager points.",
    ]
    for i, rule in enumerate(rules):
        if i == 0:
            rtf.paragraphs[0].text = rule
            rtf.paragraphs[0].space_after = Pt(8)
            rtf.paragraphs[0].runs[0].font.size = Pt(18)
        else:
            pp = rtf.add_paragraph()
            pp.text = rule
            pp.space_after = Pt(8)
            pp.runs[0].font.size = Pt(18)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
