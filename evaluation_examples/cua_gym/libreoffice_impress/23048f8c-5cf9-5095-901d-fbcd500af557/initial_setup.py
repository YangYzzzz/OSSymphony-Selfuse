"""
Initial Setup: Product 3D presentation with 7 slides, image on slide 5, no animations.
Task ID: impress_ma_084
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Generate a simple product-render-style image using PIL
from PIL import Image, ImageDraw, ImageFont

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_084'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
IMG_PATH = f'{WORKDIR}/product_render.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_product_image():
    """Create a simple 3D-style product render placeholder image."""
    img = Image.new('RGB', (600, 500), color=(240, 240, 245))
    draw = ImageDraw.Draw(img)

    # Draw a 3D-ish box to represent a product
    # Front face
    draw.polygon([(180, 150), (420, 150), (420, 380), (180, 380)],
                 fill=(50, 120, 200), outline=(30, 80, 160))
    # Top face (3D effect)
    draw.polygon([(180, 150), (230, 100), (470, 100), (420, 150)],
                 fill=(80, 150, 230), outline=(30, 80, 160))
    # Right face (3D effect)
    draw.polygon([(420, 150), (470, 100), (470, 330), (420, 380)],
                 fill=(40, 100, 180), outline=(30, 80, 160))

    # Add text label on front face
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 28)
        font_small = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 16)
    except (IOError, OSError):
        font = ImageFont.load_default()
        font_small = font

    draw.text((220, 240), "ProX 3000", fill=(255, 255, 255), font=font)
    draw.text((230, 280), "Smart Device", fill=(200, 220, 255), font=font_small)

    # Subtle shadow
    draw.ellipse([(160, 390), (460, 420)], fill=(200, 200, 205))

    img.save(IMG_PATH)
    print(f'Product image created: {IMG_PATH}')


def create_initial():
    create_product_image()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "ProX 3000 Launch"
    slide1.placeholders[1].text = "Innovative Smart Device for the Modern Home"
    bg1 = slide1.background.fill
    bg1.solid()
    bg1.fore_color.rgb = RGBColor(0x0A, 0x1F, 0x44)
    for para in slide1.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.size = Pt(40)
    for para in slide1.placeholders[1].text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)
            run.font.size = Pt(20)

    # --- Slide 2: Market Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Market Overview"
    p.runs[0].font.size = Pt(32)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(0x0A, 0x1F, 0x44)

    bullets_data = [
        "Global smart home market projected to reach $338B by 2027",
        "Year-over-year growth of 24.3% in connected device segment",
        "Consumer demand shifting toward integrated ecosystem solutions",
        "Key competitors: NexGen Home, SmartPulse, EcoLink Systems",
    ]
    for text in bullets_data:
        bp = tf.add_paragraph()
        bp.text = text
        bp.level = 0
        bp.space_before = Pt(8)
        for run in bp.runs:
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 3: Key Features ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Key Features"
    p3.runs[0].font.size = Pt(32)
    p3.runs[0].font.bold = True
    p3.runs[0].font.color.rgb = RGBColor(0x0A, 0x1F, 0x44)

    features = [
        ("Voice-Activated Control", "Supports Alexa, Google Home, and Siri integration"),
        ("Energy Monitoring", "Real-time power consumption tracking with AI optimization"),
        ("Modular Design", "Snap-on accessories for lighting, audio, and climate control"),
        ("Security Suite", "End-to-end encryption with biometric authentication"),
    ]
    for title, desc in features:
        fp = tf3.add_paragraph()
        fp.text = title
        fp.space_before = Pt(12)
        for run in fp.runs:
            run.font.size = Pt(20)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1A, 0x5C, 0xB0)
        dp = tf3.add_paragraph()
        dp.text = desc
        dp.level = 1
        for run in dp.runs:
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # --- Slide 4: Technical Specifications ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Technical Specifications"
    p4.runs[0].font.size = Pt(32)
    p4.runs[0].font.bold = True
    p4.runs[0].font.color.rgb = RGBColor(0x0A, 0x1F, 0x44)

    table_shape = slide4.shapes.add_table(6, 2, Inches(1.5), Inches(1.8), Inches(8), Inches(4))
    table = table_shape.table
    table.columns[0].width = Inches(3)
    table.columns[1].width = Inches(5)

    specs = [
        ("Specification", "Details"),
        ("Processor", "Quad-core ARM Cortex-A78, 2.4 GHz"),
        ("Connectivity", "Wi-Fi 6E, Bluetooth 5.3, Zigbee, Thread"),
        ("Power Supply", "USB-C, 15W; Battery backup 4,000 mAh"),
        ("Dimensions", "142mm x 142mm x 58mm"),
        ("Weight", "320g (base unit)"),
    ]
    for r, (k, v) in enumerate(specs):
        cell_k = table.cell(r, 0)
        cell_v = table.cell(r, 1)
        cell_k.text = k
        cell_v.text = v
        for cell in (cell_k, cell_v):
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(14)
                    if r == 0:
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    else:
                        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 5: 3D Product Model (image in center, NO animations) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox5 = slide5.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(10), Inches(1))
    tf5 = txBox5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "3D Product Model"
    p5.runs[0].font.size = Pt(32)
    p5.runs[0].font.bold = True
    p5.runs[0].font.color.rgb = RGBColor(0x0A, 0x1F, 0x44)

    # Center the image on slide 5
    img_w = Inches(5.5)
    img_h = Inches(4.6)
    left = (prs.slide_width - img_w) // 2
    top = Inches(1.6)
    pic = slide5.shapes.add_picture(IMG_PATH, left, top, img_w, img_h)
    # Give the picture a descriptive name
    pic.name = "Product 3D Render"

    # --- Slide 6: Pricing Strategy ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox6 = slide6.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1))
    tf6 = txBox6.text_frame
    p6 = tf6.paragraphs[0]
    p6.text = "Pricing Strategy"
    p6.runs[0].font.size = Pt(32)
    p6.runs[0].font.bold = True
    p6.runs[0].font.color.rgb = RGBColor(0x0A, 0x1F, 0x44)

    pricing_items = [
        "Base Unit (ProX 3000): $249.99",
        "Lighting Module Add-on: $59.99",
        "Audio Module Add-on: $79.99",
        "Climate Module Add-on: $89.99",
        "Complete Bundle (all modules): $399.99 (save $79.96)",
        "Enterprise licensing: Contact sales for volume pricing",
    ]
    for text in pricing_items:
        pp = tf6.add_paragraph()
        pp.text = text
        pp.space_before = Pt(8)
        for run in pp.runs:
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 7: Launch Timeline ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox7 = slide7.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1))
    tf7 = txBox7.text_frame
    p7 = tf7.paragraphs[0]
    p7.text = "Launch Timeline"
    p7.runs[0].font.size = Pt(32)
    p7.runs[0].font.bold = True
    p7.runs[0].font.color.rgb = RGBColor(0x0A, 0x1F, 0x44)

    timeline = [
        ("Q2 2026", "Beta testing with 500 selected households"),
        ("Q3 2026", "Manufacturing ramp-up and retail partnerships"),
        ("October 2026", "Official launch at CES Asia"),
        ("Q4 2026", "Global rollout - North America and Europe"),
        ("Q1 2027", "Asia-Pacific expansion and accessory ecosystem launch"),
    ]
    for date, milestone in timeline:
        tp = tf7.add_paragraph()
        tp.space_before = Pt(10)
        r1 = tp.add_run()
        r1.text = f"{date}: "
        r1.font.size = Pt(18)
        r1.font.bold = True
        r1.font.color.rgb = RGBColor(0x1A, 0x5C, 0xB0)
        r2 = tp.add_run()
        r2.text = milestone
        r2.font.size = Pt(18)
        r2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
