"""
Initial Setup: Safety training presentation with bulleted lists
Task ID: impstruct_029
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impstruct_029'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_bullet_paragraph(tf, text, level=0, font_size=Pt(18), is_first=False):
    """Add a bulleted paragraph to a text frame."""
    if is_first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.space_after = Pt(6)
    for run in p.runs:
        run.font.size = font_size
        run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    return p


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Workplace Safety Training 2025"
    slide1.placeholders[1].text = "Comprehensive Safety Protocols & Best Practices\nHR Department — Quarterly Review"

    # --- Slide 2: General Safety Guidelines ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "General Safety Guidelines"
    tf2 = slide2.placeholders[1].text_frame
    tf2.clear()
    bullets_s2 = [
        "Always wear appropriate personal protective equipment (PPE) in designated areas",
        "Report any unsafe conditions or near-miss incidents to your supervisor immediately",
        "Keep walkways and emergency exits clear of obstructions at all times",
        "Follow all posted safety signs and floor markings throughout the facility",
        "Attend mandatory safety briefings at the start of each shift",
        "Know the location of fire extinguishers and first aid kits on your floor",
    ]
    for i, text in enumerate(bullets_s2):
        add_bullet_paragraph(tf2, text, level=0, is_first=(i == 0))

    # --- Slide 3: Emergency Procedures ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Emergency Procedures"
    tf3 = slide3.placeholders[1].text_frame
    tf3.clear()
    bullets_s3 = [
        "Evacuate calmly using the nearest marked exit when the alarm sounds",
        "Assemble at designated muster points — Building A: North Parking Lot",
        "Do not use elevators during fire emergencies under any circumstances",
        "Call 911 for medical emergencies, then notify the front desk at ext. 200",
        "Designated floor wardens must conduct headcounts within 5 minutes",
        "Review evacuation maps posted near each stairwell entrance monthly",
    ]
    for i, text in enumerate(bullets_s3):
        add_bullet_paragraph(tf3, text, level=0, is_first=(i == 0))

    # --- Slide 4: Hazardous Materials Handling ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Hazardous Materials Handling"
    tf4 = slide4.placeholders[1].text_frame
    tf4.clear()
    bullets_s4 = [
        "Consult Safety Data Sheets (SDS) before handling any chemical substance",
        "Store flammable liquids in approved cabinets away from ignition sources",
        "Use chemical fume hoods when working with volatile or toxic compounds",
        "Label all secondary containers with the chemical name and hazard warnings",
        "Dispose of hazardous waste only through the Environmental Health & Safety office",
        "Spill kits are located in Lab Rooms 102, 205, and the loading dock area",
    ]
    for i, text in enumerate(bullets_s4):
        add_bullet_paragraph(tf4, text, level=0, is_first=(i == 0))

    # --- Slide 5: Ergonomics & Workplace Wellness ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Ergonomics & Workplace Wellness"
    tf5 = slide5.placeholders[1].text_frame
    tf5.clear()
    bullets_s5 = [
        "Adjust your chair height so feet rest flat on the floor with knees at 90 degrees",
        "Position your monitor at arm's length with the top of the screen at eye level",
        "Take a 5-minute stretch break every 60 minutes to reduce repetitive strain",
        "Use ergonomic keyboards and wrist rests to minimize carpal tunnel risk",
        "Report persistent discomfort to HR for a workstation assessment",
        "Standing desk options are available — submit a request through the facilities portal",
    ]
    for i, text in enumerate(bullets_s5):
        add_bullet_paragraph(tf5, text, level=0, is_first=(i == 0))

    # --- Slide 6: Summary & Next Steps ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Summary & Next Steps"
    tf6 = slide6.placeholders[1].text_frame
    tf6.clear()
    summary_items = [
        "Complete the online safety quiz by March 28, 2025",
        "Schedule your annual physical with Occupational Health Services",
        "Review updated emergency contact lists with your team lead",
        "Questions? Contact the Safety Office at safety@company.com or ext. 450",
    ]
    for i, text in enumerate(summary_items):
        add_bullet_paragraph(tf6, text, level=0, is_first=(i == 0))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
