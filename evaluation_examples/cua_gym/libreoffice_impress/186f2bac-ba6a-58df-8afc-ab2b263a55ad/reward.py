"""
Reward Script: Insert a line chart on slide 5 with monthly website traffic data
Task ID: impress_rp_022
Domain: libreoffice_impress
Scoring:
  Component 1: Chart exists on slide 5 and is a line chart (0.20)
  Component 2: Chart title is 'Monthly Website Traffic' (0.15)
  Component 3: Correct data values for 6 months (0.25)
  Component 4: Correct categories Jan-Jun (0.15)
  Component 5: Green line color #27AE60 with ~2pt width (0.15)
  Component 6: Circular markers on the line (0.10)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_022'


def persist_app_state():
    """Save any unsaved LibreOffice state before verification."""
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        import time
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # Slide 5 (0-indexed)

    # Find chart shape on slide 5
    chart_shape = None
    for shape in slide.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Chart exists on slide 5 and is a line chart type (0.20 points)
    try:
        if chart_shape is None:
            print("FAIL: Component 1 -- No chart found on slide 5")
        else:
            chart = chart_shape.chart
            chart_type = chart.chart_type
            # LINE = 4, LINE_MARKERS = 65, LINE_STACKED = 63, etc.
            # Accept line types: 4 (LINE), 65 (LINE_MARKERS), 63 (LINE_STACKED), 66 (LINE_MARKERS_STACKED)
            line_types = {4, 63, 64, 65, 66, 67, 68}
            if chart_type in line_types:
                print(f"PASS: Component 1 -- Line chart found on slide 5 (type={chart_type}) (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 1 -- Chart exists but is not a line chart (type={chart_type})")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # If no chart found, remaining components cannot be checked
    if chart_shape is None:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    chart = chart_shape.chart

    # Component 2: Chart title is 'Monthly Website Traffic' (0.15 points)
    try:
        if chart.has_title:
            title_text = chart.chart_title.text_frame.text.strip()
            if title_text == 'Monthly Website Traffic':
                print(f"PASS: Component 2 -- Chart title is 'Monthly Website Traffic' (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 2 -- Chart title is '{title_text}', expected 'Monthly Website Traffic'")
        else:
            print("FAIL: Component 2 -- Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Correct data values for 6 months (0.25 points)
    expected_values = [12000.0, 15000.0, 18000.0, 22000.0, 19000.0, 25000.0]
    try:
        if len(chart.series) >= 1:
            actual_values = list(chart.series[0].values)
            if len(actual_values) == 6:
                match_count = sum(
                    1 for a, e in zip(actual_values, expected_values)
                    if abs(float(a) - e) < 1.0
                )
                if match_count == 6:
                    print(f"PASS: Component 3 -- All 6 data values match (0.25 pts)")
                    total_score += 0.25
                else:
                    print(f"FAIL: Component 3 -- {match_count}/6 values match. Actual: {actual_values}")
            else:
                print(f"FAIL: Component 3 -- Expected 6 data points, found {len(actual_values)}")
        else:
            print("FAIL: Component 3 -- No series found in chart")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Correct categories Jan-Jun (0.15 points)
    expected_cats = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']
    try:
        plot = chart.plots[0]
        actual_cats = list(plot.categories)
        if len(actual_cats) == 6:
            match_count = sum(
                1 for a, e in zip(actual_cats, expected_cats)
                if str(a).strip().lower() == e.lower()
            )
            if match_count == 6:
                print(f"PASS: Component 4 -- All categories match Jan-Jun (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 4 -- {match_count}/6 categories match. Actual: {actual_cats}")
        else:
            print(f"FAIL: Component 4 -- Expected 6 categories, found {len(actual_cats)}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Green line color #27AE60 with ~2pt width (0.15 points)
    try:
        from lxml import etree
        ser_elem = chart.series[0]._element
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        # Check line color
        ln = ser_elem.find(f'.//{{{ns_a}}}ln')
        color_ok = False
        width_ok = False

        if ln is not None:
            srgb = ln.find(f'.//{{{ns_a}}}srgbClr')
            if srgb is not None:
                color_val = srgb.get('val', '').upper()
                if color_val == '27AE60':
                    color_ok = True
                else:
                    print(f"  INFO: Line color is #{color_val}, expected #27AE60")

            # Check line width: 2pt = 25400 EMU (tolerance: 20000-32000 for ~1.5-2.5pt)
            w_attr = ln.get('w')
            if w_attr is not None:
                w_val = int(w_attr)
                if 20000 <= w_val <= 32000:
                    width_ok = True
                else:
                    print(f"  INFO: Line width is {w_val} EMU ({w_val/12700:.1f}pt), expected ~25400 (2pt)")

        sub_score = 0.0
        if color_ok and width_ok:
            sub_score = 0.15
            print(f"PASS: Component 5 -- Line color #27AE60 and ~2pt width (0.15 pts)")
        elif color_ok:
            sub_score = 0.10
            print(f"PARTIAL: Component 5 -- Line color correct but width incorrect (0.10 pts)")
        elif width_ok:
            sub_score = 0.05
            print(f"PARTIAL: Component 5 -- Line width correct but color incorrect (0.05 pts)")
        else:
            print(f"FAIL: Component 5 -- Neither line color nor width match")
        total_score += sub_score
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # Component 6: Circular markers on the line (0.10 points)
    try:
        from lxml import etree
        ser_elem = chart.series[0]._element
        ns_c = 'http://schemas.openxmlformats.org/drawingml/2006/chart'

        marker = ser_elem.find(f'{{{ns_c}}}marker')
        if marker is not None:
            symbol = marker.find(f'{{{ns_c}}}symbol')
            if symbol is not None:
                symbol_val = symbol.get('val', '')
                if symbol_val == 'circle':
                    print(f"PASS: Component 6 -- Circular markers found (0.10 pts)")
                    total_score += 0.10
                else:
                    print(f"FAIL: Component 6 -- Marker symbol is '{symbol_val}', expected 'circle'")
            else:
                print("FAIL: Component 6 -- Marker element has no symbol")
        else:
            print("FAIL: Component 6 -- No marker element found on series")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
