"""
Initial Setup: Create Digital_Report.pptx with 8 slides, slide 5 titled 'Web Analytics' with empty content
Task ID: impress_rp_022
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_022'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_body_slide(prs, title_text, body_lines):
    """Add a slide with title and bullet-point body text."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Digital Marketing Report"
    slide1.placeholders[1].text = "Q2 2025 Performance Overview\nPrepared by Analytics Team"

    # --- Slide 2: Executive Summary ---
    add_title_body_slide(prs, "Executive Summary", [
        "Overall digital engagement increased 23% quarter-over-quarter",
        "Email campaigns achieved a 4.2% average open rate",
        "Social media following grew by 18,500 new followers",
        "Paid advertising ROAS improved to 3.8x from 2.9x",
        "Mobile traffic now accounts for 67% of all sessions",
    ])

    # --- Slide 3: Campaign Performance ---
    add_title_body_slide(prs, "Campaign Performance", [
        "Spring Launch Campaign: 2.1M impressions, 45K clicks (CTR 2.1%)",
        "Newsletter Redesign: Open rate up from 3.1% to 4.8%",
        "Influencer Partnership: 890K reach, 12K new signups",
        "Retargeting Ads: Conversion rate 6.3%, CPA $12.40",
        "Content Series: 340K video views, 28% completion rate",
    ])

    # --- Slide 4: Social Media Metrics ---
    add_title_body_slide(prs, "Social Media Metrics", [
        "Instagram: 52,400 followers (+14%), Engagement rate 3.7%",
        "LinkedIn: 28,100 followers (+22%), Post reach avg 8,200",
        "Twitter/X: 19,800 followers (+8%), Avg impressions 15K",
        "TikTok: 41,300 followers (+35%), Avg views 45K per video",
        "Facebook: 34,600 followers (+5%), Group members 12,400",
    ])

    # --- Slide 5: Web Analytics (EMPTY CONTENT - task target) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide5.shapes.title.text = "Web Analytics"
    # Leave content placeholder empty - clear it
    tf5 = slide5.placeholders[1].text_frame
    tf5.clear()
    tf5.paragraphs[0].text = ""

    # --- Slide 6: Conversion Funnel ---
    add_title_body_slide(prs, "Conversion Funnel Analysis", [
        "Awareness Stage: 485,000 unique visitors per month",
        "Interest Stage: 72,000 page views on product pages (14.8%)",
        "Consideration Stage: 18,400 demo requests (25.6% of interest)",
        "Decision Stage: 4,200 trial signups (22.8% of consideration)",
        "Purchase Stage: 1,680 conversions (40% of trials)",
    ])

    # --- Slide 7: Budget Allocation ---
    add_title_body_slide(prs, "Budget Allocation", [
        "Paid Search (Google Ads): $45,000 (30% of budget)",
        "Social Media Advertising: $37,500 (25% of budget)",
        "Content Marketing & SEO: $30,000 (20% of budget)",
        "Email Marketing Platform: $15,000 (10% of budget)",
        "Influencer Partnerships: $22,500 (15% of budget)",
    ])

    # --- Slide 8: Next Quarter Priorities ---
    add_title_body_slide(prs, "Q3 2025 Priorities", [
        "Launch AI-powered chatbot for website visitor engagement",
        "Expand TikTok content strategy with 3x weekly posting",
        "Implement advanced attribution modeling across channels",
        "Redesign landing pages for mobile-first experience",
        "Increase budget allocation to high-ROAS channels by 15%",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
