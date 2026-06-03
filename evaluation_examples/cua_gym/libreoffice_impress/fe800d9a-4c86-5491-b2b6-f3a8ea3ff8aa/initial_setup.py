"""
Initial Setup: Decision Tree branching presentation with 12 slides
Task ID: impress_gf2_041
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_041'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_text(slide, title_text, body_text=None):
    """Add a title textbox and optional body text to a blank slide."""
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(8.4), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6D)

    if body_text:
        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(4.5))
        btf = body_box.text_frame
        btf.word_wrap = True
        for i, line in enumerate(body_text if isinstance(body_text, list) else [body_text]):
            if i == 0:
                bp = btf.paragraphs[0]
            else:
                bp = btf.add_paragraph()
            brun = bp.add_run()
            brun.text = line
            brun.font.size = Pt(16)
            brun.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_text(slide1, "Strategic Market Expansion: Decision Tree Analysis", [
        "Prepared by the Strategy & Analytics Division",
        "Q2 2025 Planning Session",
        "",
        "This interactive presentation guides stakeholders through",
        "key decision points for the proposed APAC market expansion."
    ])

    # Slide 2: Introduction / Overview
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_text(slide2, "How to Navigate This Presentation", [
        "This deck uses a branching decision tree structure.",
        "",
        "At each decision point, you will see navigation buttons.",
        "Click the button that matches your assessment to follow",
        "the corresponding analysis path.",
        "",
        "Key decision points appear on slides 3, 7, and 11."
    ])

    # Slide 3: Decision Question (space at bottom for buttons, but NO buttons yet)
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_text(slide3, "Decision Point #1: Market Readiness", [
        "Based on the preliminary market analysis, our research team",
        "has identified two possible interpretations of the data:",
        "",
        "Option A: The APAC consumer electronics market shows strong",
        "growth indicators and favorable regulatory conditions.",
        "",
        "Option B: Currency volatility and supply chain risks in the",
        "region suggest a more cautious approach is warranted.",
        "",
        "Does the data support immediate market entry?"
    ])
    # Note: NO buttons here - that's the task for the agent

    # Slide 4: Yes outcome path start
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_text(slide4, "Path A: Accelerated Market Entry", [
        "Given the positive market indicators, we recommend the",
        "following accelerated entry timeline:",
        "",
        "Phase 1 (Q3 2025): Establish regional partnerships",
        "  - Partner with Toshiba Electronics for distribution",
        "  - Secure warehouse space in Singapore and Tokyo",
        "",
        "Phase 2 (Q4 2025): Soft launch in 3 pilot markets",
        "  - Singapore, South Korea, and Japan",
        "  - Initial inventory: 50,000 units per market"
    ])

    # Slide 5: Continuation of Yes path
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_text(slide5, "Path A: Financial Projections", [
        "Revenue Forecast (Accelerated Entry):",
        "",
        "Year 1: $12.4M (conservative) - $18.7M (optimistic)",
        "Year 2: $28.1M (conservative) - $42.3M (optimistic)",
        "Year 3: $45.6M (conservative) - $68.9M (optimistic)",
        "",
        "Break-even expected within 14-18 months.",
        "Required initial investment: $8.2M",
        "Expected ROI at 36 months: 285% - 430%"
    ])

    # Slide 6: More Yes path details
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_text(slide6, "Path A: Risk Mitigation Measures", [
        "To address residual risks in the accelerated timeline:",
        "",
        "1. Currency hedging contracts with Goldman Sachs",
        "   - Lock in USD/JPY and USD/SGD rates for 12 months",
        "",
        "2. Dual-supplier strategy for critical components",
        "   - Primary: Samsung Display Co.",
        "   - Secondary: LG Innotek",
        "",
        "3. Regional insurance coverage via AIG Asia Pacific"
    ])

    # Slide 7: Last slide of Yes branch
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_text(slide7, "Path A: Conclusion & Next Steps", [
        "The accelerated entry strategy offers significant upside",
        "potential with manageable risk exposure.",
        "",
        "Recommended next steps:",
        "  - Board approval for $8.2M capital allocation",
        "  - Legal review of partnership agreements by June 15",
        "  - Hire Regional Director (APAC) by end of Q2",
        "",
        "Timeline to first revenue: 6-8 months from approval."
    ])
    # Note: NO "Go to Summary" button here - that's the task

    # Slide 8: No outcome path start
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_text(slide8, "Path B: Cautious Evaluation Period", [
        "Given the identified risks, a phased evaluation",
        "approach is recommended before committing resources:",
        "",
        "Phase 1 (Q3-Q4 2025): Extended market research",
        "  - Commission McKinsey for deep-dive analysis",
        "  - Monitor currency trends and trade policy changes",
        "",
        "Phase 2 (Q1 2026): Limited pilot program",
        "  - Single market entry (Singapore only)",
        "  - Minimal inventory commitment: 10,000 units"
    ])

    # Slide 9: Continuation of No path
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_text(slide9, "Path B: Risk Assessment Details", [
        "Key risk factors requiring monitoring:",
        "",
        "1. Currency Risk: JPY/USD volatility at 18-month high",
        "   - Potential impact: 8-15% margin erosion",
        "",
        "2. Supply Chain: Semiconductor shortage ongoing",
        "   - Lead times extended to 26 weeks (was 12 weeks)",
        "",
        "3. Regulatory: New data privacy laws in South Korea",
        "   - Compliance costs estimated at $1.2M annually",
        "",
        "4. Competition: Xiaomi expanding aggressively in SE Asia"
    ])

    # Slide 10: More No path
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_text(slide10, "Path B: Conservative Financial Outlook", [
        "Revenue Forecast (Cautious Approach):",
        "",
        "Year 1: $3.8M (pilot market only)",
        "Year 2: $14.2M (gradual expansion to 3 markets)",
        "Year 3: $32.7M (full regional presence)",
        "",
        "Break-even expected within 22-28 months.",
        "Required initial investment: $3.5M",
        "Expected ROI at 36 months: 165% - 240%",
        "",
        "Lower upside but significantly reduced downside risk."
    ])

    # Slide 11: Summary
    slide11 = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_text(slide11, "Executive Summary: Both Paths Compared", [
        "Path A (Accelerated):        Path B (Cautious):",
        "  Investment: $8.2M            Investment: $3.5M",
        "  Break-even: 14-18 mo         Break-even: 22-28 mo",
        "  3-Year ROI: 285-430%         3-Year ROI: 165-240%",
        "  Risk Level: Moderate          Risk Level: Low",
        "",
        "Both strategies are viable. The choice depends on the",
        "board's risk appetite and available capital allocation.",
        "",
        "Recommendation: Path A if capital is available and board",
        "approves the risk profile. Path B otherwise."
    ])
    # Note: NO "Start Over" or "End Presentation" buttons - that's the task

    # Slide 12: End / Thank You
    slide12 = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_text(slide12, "Thank You", [
        "Strategic Market Expansion: Decision Tree Analysis",
        "",
        "Prepared by: Sarah Chen, VP Strategy & Analytics",
        "Reviewed by: Marcus Johnson, CFO",
        "",
        "Contact: strategy@acmecorp.com",
        "Date: March 2025",
        "",
        "Confidential - Internal Use Only"
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
