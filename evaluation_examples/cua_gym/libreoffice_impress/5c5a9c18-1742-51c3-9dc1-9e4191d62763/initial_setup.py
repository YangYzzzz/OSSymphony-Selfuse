"""
Initial Setup: timeline_presentation.pptx with 5 slides, slide 5 has timeline with M1-M5 milestone shapes, NO animations.
Task ID: impress_anim_060
Domain: libreoffice_impress
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn as pptx_qn

WORKDIR = '/home/user/Desktop'
TASK_ID = 'impress_anim_060'
OUTPUT = f'{WORKDIR}/timeline_presentation.pptx'


def set_text_anchor_center(text_frame):
    """Set vertical text anchor to center using XML."""
    # The attribute is 'anchor' on txBody
    text_frame._txBody.set('anchor', 'ctr')


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_layouts = prs.slide_layouts

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(slide_layouts[0])
    slide1.shapes.title.text = "Project Timeline"
    slide1.placeholders[1].text = "Q1 2025 Roadmap Overview"
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x5C)
    for shape in slide1.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "This roadmap outlines five major milestones for Q1 2025."
    tf2.add_paragraph().text = "Each milestone represents a critical deliverable."
    tf2.add_paragraph().text = "Teams are aligned on delivery dates and dependencies."
    tf2.add_paragraph().text = "Budget approved: $2.4M across all initiatives."

    # --- Slide 3: Team Overview ---
    slide3 = prs.slides.add_slide(slide_layouts[1])
    slide3.shapes.title.text = "Team Overview"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Engineering: 12 developers, 3 QA engineers"
    tf3.add_paragraph().text = "Product: 2 PMs, 1 UX designer"
    tf3.add_paragraph().text = "Marketing: 4 specialists"
    tf3.add_paragraph().text = "Operations: 3 support staff"

    # --- Slide 4: Resource Allocation ---
    slide4 = prs.slides.add_slide(slide_layouts[1])
    slide4.shapes.title.text = "Resource Allocation"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "M1 - Discovery & Planning: 3 weeks, $400K"
    tf4.add_paragraph().text = "M2 - Architecture & Design: 4 weeks, $550K"
    tf4.add_paragraph().text = "M3 - Core Development: 6 weeks, $720K"
    tf4.add_paragraph().text = "M4 - Testing & QA: 3 weeks, $380K"
    tf4.add_paragraph().text = "M5 - Launch & Rollout: 2 weeks, $350K"

    # --- Slide 5: Timeline (Main Slide) ---
    slide5 = prs.slides.add_slide(slide_layouts[5])  # Blank layout
    slide5.background.fill.solid()
    slide5.background.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)

    # Title text box at top
    title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_run = title_p.add_run()
    title_run.text = "Q1 2025 Project Timeline"
    title_run.font.bold = True
    title_run.font.size = Pt(28)
    title_run.font.color.rgb = RGBColor(0x1A, 0x3A, 0x5C)

    # Horizontal timeline bar (thin rectangle)
    bar = slide5.shapes.add_shape(
        1,  # rounded rect
        Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.18)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x4A, 0x90, 0xD9)
    bar.line.color.rgb = RGBColor(0x4A, 0x90, 0xD9)
    bar.name = "TimelineBar"

    # Define milestone positions evenly across the timeline
    milestone_data = [
        ("M1", "Discovery\n& Planning", Inches(1.2), RGBColor(0x27, 0xAE, 0x60)),
        ("M2", "Architecture\n& Design", Inches(3.5), RGBColor(0x29, 0x80, 0xB9)),
        ("M3", "Core\nDevelopment", Inches(5.8), RGBColor(0x8E, 0x44, 0xAD)),
        ("M4", "Testing\n& QA", Inches(8.1), RGBColor(0xE6, 0x7E, 0x22)),
        ("M5", "Launch\n& Rollout", Inches(10.4), RGBColor(0xC0, 0x39, 0x2B)),
    ]

    circle_size = Inches(1.1)
    circle_top = Inches(3.05)

    for label, desc_text, left_pos, color in milestone_data:
        # Use oval shape (type 9)
        circle = slide5.shapes.add_shape(
            9,  # oval
            left_pos, circle_top, circle_size, circle_size
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        circle.line.width = Pt(2.5)
        circle.name = label  # NAME the shape M1, M2, M3, M4, M5

        # Label inside circle
        circle_tf = circle.text_frame
        circle_tf.word_wrap = False
        circle_p = circle_tf.paragraphs[0]
        circle_p.alignment = PP_ALIGN.CENTER
        circle_run = circle_p.add_run()
        circle_run.text = label
        circle_run.font.bold = True
        circle_run.font.size = Pt(16)
        circle_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        set_text_anchor_center(circle_tf)

        # Description text box below the circle
        desc_box = slide5.shapes.add_textbox(
            left_pos - Inches(0.1), circle_top + circle_size + Inches(0.15),
            circle_size + Inches(0.2), Inches(0.9)
        )
        desc_tf = desc_box.text_frame
        desc_tf.word_wrap = True
        desc_p = desc_tf.paragraphs[0]
        desc_p.alignment = PP_ALIGN.CENTER
        desc_run = desc_p.add_run()
        desc_run.text = desc_text
        desc_run.font.size = Pt(10)
        desc_run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # Date labels on timeline bar
    dates = ["Jan 6", "Jan 27", "Feb 24", "Mar 17", "Mar 31"]
    date_lefts = [Inches(1.1), Inches(3.4), Inches(5.7), Inches(8.0), Inches(10.3)]
    for date_text, date_left in zip(dates, date_lefts):
        date_box = slide5.shapes.add_textbox(
            date_left, Inches(3.85), Inches(1.3), Inches(0.35)
        )
        date_tf = date_box.text_frame
        date_p = date_tf.paragraphs[0]
        date_p.alignment = PP_ALIGN.CENTER
        date_run = date_p.add_run()
        date_run.text = date_text
        date_run.font.size = Pt(9)
        date_run.font.color.rgb = RGBColor(0x77, 0x77, 0x77)

    # Subtitle at bottom
    subtitle_box = slide5.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4))
    subtitle_tf = subtitle_box.text_frame
    subtitle_p = subtitle_tf.paragraphs[0]
    subtitle_p.alignment = PP_ALIGN.CENTER
    subtitle_run = subtitle_p.add_run()
    subtitle_run.text = "All dates subject to change | Engineering Division | Confidential"
    subtitle_run.font.size = Pt(9)
    subtitle_run.font.italic = True
    subtitle_run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    # Ensure Desktop directory exists (on VM)
    os.makedirs(WORKDIR, exist_ok=True)
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Slide 5 has 5 milestone shapes (M1-M5) with NO animations.')


create_initial()
