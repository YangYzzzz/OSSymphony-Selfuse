"""
Initial Setup: Create a 5-slide project status presentation with black text.
Task ID: impstruct_019
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impstruct_019'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

BLACK = RGBColor(0x00, 0x00, 0x00)


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    for run in slide.shapes.title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = BLACK
        run.font.size = Pt(36)
        run.font.bold = True
    slide.placeholders[1].text = subtitle_text
    for run in slide.placeholders[1].text_frame.paragraphs[0].runs:
        run.font.color.rgb = BLACK
        run.font.size = Pt(20)
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    for run in slide.shapes.title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = BLACK
        run.font.size = Pt(28)
        run.font.bold = True

    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.space_after = Pt(6)
        for run in p.runs:
            run.font.color.rgb = BLACK
            run.font.size = Pt(18)
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    add_title_slide(
        prs,
        "Project Phoenix - Status Report",
        "Quarterly Review | April 2026 | Engineering Division"
    )

    # --- Slide 2: Executive Summary ---
    add_content_slide(prs, "Executive Summary", [
        "Overall project health: On Track",
        "Budget utilization at 67% ($2.4M of $3.6M allocated)",
        "Team expanded from 12 to 18 engineers this quarter",
        "Three critical modules completed ahead of schedule",
        "Client satisfaction score: 4.7/5.0 based on latest survey",
        "Next major delivery milestone: June 15, 2026",
    ])

    # --- Slide 3: Milestones (target slide for the task) ---
    add_content_slide(prs, "Milestones", [
        "Q1 2026: Requirements gathering and architecture design - Complete",
        "Q2 2026: Core API development and database migration - In Progress",
        "Q3 2026: Frontend redesign and integration testing - Planned",
        "Q4 2026: Performance optimization and security audit - Planned",
        "January 2027: Beta release to select enterprise clients",
        "March 2027: General availability launch across all regions",
    ])

    # --- Slide 4: Risk Assessment ---
    add_content_slide(prs, "Risk Assessment", [
        "Supply chain delays may impact hardware provisioning timeline",
        "Two senior developers transitioning to advisory roles in Q3",
        "Regulatory compliance review pending for EU market entry",
        "Third-party API deprecation requires migration by August 2026",
        "Mitigation plans documented and approved by steering committee",
    ])

    # --- Slide 5: Next Steps ---
    add_content_slide(prs, "Next Steps", [
        "Complete sprint 14 deliverables by April 18, 2026",
        "Schedule architecture review with external consultants",
        "Finalize vendor contracts for cloud infrastructure expansion",
        "Begin user acceptance testing with pilot group of 50 users",
        "Prepare mid-year budget revision proposal for leadership",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
