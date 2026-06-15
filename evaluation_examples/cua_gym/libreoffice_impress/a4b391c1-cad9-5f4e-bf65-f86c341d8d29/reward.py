"""
Reward Script: Change font color of all text on slide 3 to dark blue (#003366)
Task ID: impstruct_019
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Title text on slide 3 has color #003366
  Component 2 (0.4): All body text runs on slide 3 have color #003366
  Component 3 (0.2): Text on other slides remains unchanged (not #003366)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impstruct_019'
TARGET_COLOR = '003366'


def get_run_color_hex(run):
    """Get the hex color string of a run's font, or None if inherited/unset."""
    try:
        if run.font.color.type is not None:
            return str(run.font.color.rgb)
    except Exception:
        pass
    return None


def get_all_nonempty_runs(slide):
    """Return all runs with non-empty text from all text-bearing shapes on a slide."""
    runs = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        runs.append(run)
    return runs


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]  # 0-indexed

    # Component 1: Title text on slide 3 has color #003366 (0.4 points)
    try:
        title_runs = []
        for shape in slide3.shapes:
            if shape.has_text_frame and shape.name and 'Title' in shape.name:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if (run.text or "").strip():
                            title_runs.append(run)

        if not title_runs:
            print("FAIL: Component 1 -- No title runs found on slide 3")
        else:
            all_title_correct = True
            for run in title_runs:
                color = get_run_color_hex(run)
                if color != TARGET_COLOR:
                    all_title_correct = False
                    print(f"FAIL: Component 1 -- Title run '{run.text[:30]}' has color {color}, expected {TARGET_COLOR}")
                    break

            if all_title_correct:
                print(f"PASS: Component 1 -- All title text on slide 3 has color #{TARGET_COLOR} (0.4 pts)")
                total_score += 0.4
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: All body text runs on slide 3 have color #003366 (0.4 points)
    try:
        body_runs = []
        for shape in slide3.shapes:
            if shape.has_text_frame and not (shape.name and 'Title' in shape.name):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if (run.text or "").strip():
                            body_runs.append(run)

        if not body_runs:
            print("FAIL: Component 2 -- No body text runs found on slide 3")
        else:
            correct_count = 0
            for run in body_runs:
                color = get_run_color_hex(run)
                if color == TARGET_COLOR:
                    correct_count += 1
                else:
                    print(f"FAIL: Component 2 -- Body run '{run.text[:30]}' has color {color}, expected {TARGET_COLOR}")

            if correct_count == len(body_runs):
                print(f"PASS: Component 2 -- All {len(body_runs)} body text runs on slide 3 have color #{TARGET_COLOR} (0.4 pts)")
                total_score += 0.4
            else:
                # Partial credit: proportional
                partial = 0.4 * (correct_count / len(body_runs))
                print(f"PARTIAL: Component 2 -- {correct_count}/{len(body_runs)} body runs correct ({partial:.2f} pts)")
                total_score += partial
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Slide 3 ALL text is #003366 AND other slides have NO #003366 text (0.2 points)
    # This compound check ensures the change was correctly scoped to slide 3 only.
    # It anchors on the task-introduced change (slide 3 being dark blue) so it fails on initial_env.
    try:
        # First: confirm slide 3 has the change (anchor condition)
        slide3_runs = get_all_nonempty_runs(slide3)
        slide3_all_correct = all(get_run_color_hex(r) == TARGET_COLOR for r in slide3_runs) if slide3_runs else False

        if not slide3_all_correct:
            print("FAIL: Component 3 -- Slide 3 text is not all #003366, so scoping check is moot")
        else:
            # Second: confirm other slides are NOT contaminated
            other_slide_indices = [i for i in range(len(prs.slides)) if i != 2]
            contaminated = False
            for idx in other_slide_indices:
                slide = prs.slides[idx]
                runs = get_all_nonempty_runs(slide)
                for run in runs:
                    color = get_run_color_hex(run)
                    if color == TARGET_COLOR:
                        contaminated = True
                        print(f"FAIL: Component 3 -- Slide {idx+1} run '{run.text[:30]}' has color #{TARGET_COLOR} (should be unchanged)")
                        break
                if contaminated:
                    break

            if not contaminated:
                print(f"PASS: Component 3 -- Slide 3 is all #{TARGET_COLOR} AND no other slides affected (0.2 pts)")
                total_score += 0.2
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
