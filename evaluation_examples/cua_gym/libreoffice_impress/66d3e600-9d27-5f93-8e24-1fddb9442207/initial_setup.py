"""
Initial Setup: Create a 12-slide Funding Journey presentation with slide 10 titled 'Funding Milestones' but no content.
Task ID: impress_exec_093
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_093'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    if subtitle_text and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.text = body_lines[0]
    for line in body_lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
    return slide


def add_title_only_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    add_title_slide(prs, "Funding Journey", "NovaTech AI - From Seed to IPO")

    # Slide 2: Company Overview
    add_content_slide(prs, "Company Overview", [
        "Founded in 2019 by Dr. Amelia Park and Raj Mehta",
        "Headquartered in San Francisco, CA",
        "AI-powered supply chain optimization platform",
        "150+ enterprise customers across 12 countries",
        "Revenue: $42M ARR (Q4 2025)",
    ])

    # Slide 3: Founding Story
    add_content_slide(prs, "The Founding Story", [
        "Started as a Stanford research project in 2018",
        "Won TechCrunch Disrupt in 2019 with working prototype",
        "First paying customer within 6 months of incorporation",
        "Bootstrap phase: $500K personal savings + angel checks",
    ])

    # Slide 4: Seed Round Details
    add_content_slide(prs, "Seed Round - $2M (2020)", [
        "Lead Investor: Sequoia Scout Fund",
        "Participation: Y Combinator (W20 batch)",
        "Use of funds: Core platform development, first 5 hires",
        "Key milestone: 10 pilot customers onboarded",
        "Valuation: $8M pre-money",
    ])

    # Slide 5: Series A Details
    add_content_slide(prs, "Series A - $10M (2021)", [
        "Lead Investor: Andreessen Horowitz",
        "Participation: Sequoia, Y Combinator Continuity",
        "Use of funds: Scale engineering team to 30, enterprise sales",
        "Key milestone: $5M ARR, Fortune 500 customer wins",
        "Valuation: $50M pre-money",
    ])

    # Slide 6: Series B Details
    add_content_slide(prs, "Series B - $35M (2022)", [
        "Lead Investor: Tiger Global Management",
        "Participation: a16z, Sequoia, Founders Fund",
        "Use of funds: International expansion (EMEA + APAC)",
        "Key milestone: $18M ARR, 80 enterprise customers",
        "Valuation: $200M pre-money",
    ])

    # Slide 7: Series C Details
    add_content_slide(prs, "Series C - $80M (2024)", [
        "Lead Investor: SoftBank Vision Fund",
        "Participation: All existing investors + Wellington Management",
        "Use of funds: AI R&D, strategic acquisitions, IPO readiness",
        "Key milestone: $42M ARR, profitable unit economics",
        "Valuation: $800M pre-money",
    ])

    # Slide 8: Key Metrics
    add_content_slide(prs, "Key Performance Metrics", [
        "Revenue Growth: 120% YoY (2023-2024)",
        "Net Revenue Retention: 145%",
        "Gross Margin: 78%",
        "CAC Payback Period: 14 months",
        "Employee Count: 320 across 4 offices",
    ])

    # Slide 9: Investor Relations
    add_content_slide(prs, "Investor Relations Summary", [
        "Total capital raised: $127M across 4 rounds",
        "Current valuation: $880M post-money",
        "Board composition: 2 founders, 3 investor seats, 2 independents",
        "Next milestone: Pre-IPO bridge round or direct listing",
        "Target IPO timeline: Q3 2026",
    ])

    # Slide 10: Funding Milestones - Title only, NO content
    add_title_only_slide(prs, "Funding Milestones")

    # Slide 11: IPO Readiness
    add_content_slide(prs, "IPO Readiness Checklist", [
        "SOX compliance audit completed (Deloitte)",
        "CFO hire: Jennifer Walsh (ex-Datadog)",
        "S-1 drafting in progress with Goldman Sachs",
        "Dual-class share structure approved by board",
        "Target exchange: NASDAQ under ticker NVTK",
    ])

    # Slide 12: Thank You
    add_title_slide(prs, "Thank You", "Contact: ir@novatech-ai.com | investors.novatech-ai.com")

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
