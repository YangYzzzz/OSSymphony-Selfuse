"""
Reward Script: Milestone tracker on slide 10 with 6 colored circles connected by a line
Task ID: impress_exec_093
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): Horizontal connecting line exists on slide 10
  Component 2 (0.30): 6 oval shapes exist with correct fill colors
  Component 3 (0.25): 6 text labels with correct milestone names
  Component 4 (0.15): Circle diameter is ~0.6 inches
  Component 5 (0.10): Blue circle (Series C) has a border/outline
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_093'

# Expected milestone data: (label, fill_color_hex)
EXPECTED_MILESTONES = [
    ('Seed ($2M)', '4CAF50'),
    ('Series A ($10M)', '4CAF50'),
    ('Series B ($35M)', '4CAF50'),
    ('Series C ($80M)', '2196F3'),
    ('Pre-IPO', 'E0E0E0'),
    ('IPO', 'E0E0E0'),
]


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Check slide count - must have at least 10 slides
    if len(prs.slides) < 10:
        print(f"CRITICAL: Expected at least 10 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[9]  # Slide 10 (0-indexed)

    # Collect shapes by type
    ovals = []
    rectangles = []
    text_boxes = []

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                if shape.auto_shape_type == 9:  # OVAL
                    ovals.append(shape)
                elif shape.auto_shape_type == 1:  # RECTANGLE
                    rectangles.append(shape)
            except Exception:
                pass
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            text_boxes.append(shape)
        # Also check for freeform/connector lines
        elif shape.shape_type in (MSO_SHAPE_TYPE.FREEFORM, 9, 10):  # connectors
            rectangles.append(shape)  # treat as potential line

    # ---------------------------------------------------------------
    # Component 1: Horizontal connecting line on slide 10 (0.20 pts)
    # A thin rectangle or line shape that connects the circles
    # ---------------------------------------------------------------
    try:
        # Look for a thin rectangle (height much smaller than width) that acts as connecting line
        line_found = False
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                try:
                    if shape.auto_shape_type == 1:  # RECTANGLE
                        w_in = shape.width / 914400
                        h_in = shape.height / 914400
                        # A connecting line is wide (>2 inches) and thin (<0.2 inches)
                        if w_in > 2.0 and h_in < 0.2:
                            line_found = True
                            print(f"PASS: Component 1 -- Horizontal line found ({w_in:.2f}x{h_in:.3f} in) (0.20 pts)")
                            total_score += 0.20
                            break
                except Exception:
                    pass
            # Also check for line/connector shape types
            elif shape.shape_type in (9, 10, MSO_SHAPE_TYPE.FREEFORM):
                w_in = shape.width / 914400
                if w_in > 2.0:
                    line_found = True
                    print(f"PASS: Component 1 -- Connecting line shape found (0.20 pts)")
                    total_score += 0.20
                    break

        if not line_found:
            print(f"FAIL: Component 1 -- No horizontal connecting line found on slide 10")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # ---------------------------------------------------------------
    # Component 2: 6 ovals with correct fill colors (0.30 pts)
    # Green (#4CAF50) x3, Blue (#2196F3) x1, Gray (#E0E0E0) x2
    # ---------------------------------------------------------------
    try:
        if len(ovals) < 6:
            print(f"FAIL: Component 2 -- Expected 6 ovals, found {len(ovals)}")
        else:
            # Sort ovals by horizontal position (left to right)
            sorted_ovals = sorted(ovals, key=lambda s: s.left)[:6]

            color_matches = 0
            expected_colors = ['4CAF50', '4CAF50', '4CAF50', '2196F3', 'E0E0E0', 'E0E0E0']

            for idx, (oval, expected_color) in enumerate(zip(sorted_ovals, expected_colors)):
                try:
                    actual_color = str(oval.fill.fore_color.rgb).upper()
                    if actual_color == expected_color.upper():
                        color_matches += 1
                    else:
                        print(f"  INFO: Oval {idx+1} color mismatch: expected {expected_color}, got {actual_color}")
                except Exception as e:
                    print(f"  INFO: Oval {idx+1} color read error: {e}")

            # Award proportional credit: 0.30 * (matches/6)
            component_score = 0.30 * (color_matches / 6)
            total_score += component_score
            if color_matches == 6:
                print(f"PASS: Component 2 -- All 6 ovals have correct fill colors (0.30 pts)")
            else:
                print(f"PARTIAL: Component 2 -- {color_matches}/6 ovals have correct colors ({component_score:.2f} pts)")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # ---------------------------------------------------------------
    # Component 3: 6 text labels with correct milestone names (0.25 pts)
    # Labels should appear below the circles
    # ---------------------------------------------------------------
    try:
        expected_labels = ['Seed ($2M)', 'Series A ($10M)', 'Series B ($35M)',
                           'Series C ($80M)', 'Pre-IPO', 'IPO']

        # Get all text from text boxes on slide 10 (excluding the title)
        # Filter text boxes that are in the lower portion of the slide (below circles)
        label_texts = []
        for tb in text_boxes:
            txt = tb.text_frame.text.strip() if tb.has_text_frame else ''
            if txt and txt != 'Funding Milestones':
                label_texts.append(txt)

        label_matches = 0
        for expected in expected_labels:
            if any(expected in t for t in label_texts):
                label_matches += 1
            else:
                print(f"  INFO: Label '{expected}' not found among text boxes")

        component_score = 0.25 * (label_matches / 6)
        total_score += component_score
        if label_matches == 6:
            print(f"PASS: Component 3 -- All 6 milestone labels found (0.25 pts)")
        else:
            print(f"PARTIAL: Component 3 -- {label_matches}/6 labels found ({component_score:.2f} pts)")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # ---------------------------------------------------------------
    # Component 4: Circle diameter ~0.6 inches (0.15 pts)
    # All ovals should be approximately 0.6in x 0.6in
    # ---------------------------------------------------------------
    try:
        if len(ovals) >= 6:
            sorted_ovals = sorted(ovals, key=lambda s: s.left)[:6]
            size_ok_count = 0
            target_emu = Inches(0.6)  # 548640 EMU

            for oval in sorted_ovals:
                w_ratio = abs(oval.width - target_emu) / target_emu
                h_ratio = abs(oval.height - target_emu) / target_emu
                if w_ratio <= 0.15 and h_ratio <= 0.15:  # 15% tolerance
                    size_ok_count += 1

            component_score = 0.15 * (size_ok_count / 6)
            total_score += component_score
            if size_ok_count == 6:
                print(f"PASS: Component 4 -- All 6 circles are ~0.6in diameter (0.15 pts)")
            else:
                print(f"PARTIAL: Component 4 -- {size_ok_count}/6 circles are correct size ({component_score:.2f} pts)")
        else:
            print(f"FAIL: Component 4 -- Not enough ovals to check size")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # ---------------------------------------------------------------
    # Component 5: Blue circle (Series C) has a border/outline (0.10 pts)
    # The 4th circle (blue) should have a visible border
    # ---------------------------------------------------------------
    try:
        if len(ovals) >= 4:
            sorted_ovals = sorted(ovals, key=lambda s: s.left)[:6]
            blue_oval = sorted_ovals[3]  # 4th circle (0-indexed: 3)

            # Verify it's the blue one first
            actual_color = str(blue_oval.fill.fore_color.rgb).upper()
            if actual_color == '2196F3':
                # Check for outline/border
                ln = blue_oval.line
                if ln.width is not None and ln.width > 0:
                    print(f"PASS: Component 5 -- Blue circle (Series C) has border, width={ln.width} EMU (0.10 pts)")
                    total_score += 0.10
                else:
                    print(f"FAIL: Component 5 -- Blue circle has no visible border")
            else:
                print(f"FAIL: Component 5 -- 4th circle is not blue (found {actual_color})")
        else:
            print(f"FAIL: Component 5 -- Not enough ovals to check border")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
