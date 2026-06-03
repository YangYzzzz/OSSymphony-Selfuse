"""
Reward Script: Resize slide thumbnail in notes page layout to ~30% of page height
Task ID: impress_ndo_026
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4) - Slide image height is ~30% of notes page height (reduced from ~37.5%)
  Component 2 (0.3) - Notes body top position moved up (closer to reduced slide image)
  Component 3 (0.3) - Notes body height expanded to fill remaining space
"""

import os
import sys

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_026'
FILE_PATH = f'{WORKDIR}/{TASK_ID}.pptx'

# Notes page height from presentation.xml notesSz cy
NOTES_PAGE_HEIGHT = 9144000  # EMU (10 inches)

# Initial values (pre-task state) for comparison
INITIAL_SLIDE_IMAGE_HEIGHT = 3429000   # ~37.5% of page
INITIAL_NOTES_BODY_TOP = 4343400
INITIAL_NOTES_BODY_HEIGHT = 4114800

# Expected golden values (post-task state)
# Slide image should be ~30% of page height = ~2743200 EMU
TARGET_SLIDE_IMAGE_HEIGHT_RATIO = 0.30  # 30% of notes page height
# Tolerance: accept 25%-35% as partial credit zone, 28%-32% as full credit
FULL_CREDIT_LOW = 0.28
FULL_CREDIT_HIGH = 0.32
PARTIAL_CREDIT_LOW = 0.20
PARTIAL_CREDIT_HIGH = 0.40


def persist_app_state(domain):
    """Try to save any unsaved LibreOffice state via Ctrl+S."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify that the notes page layout has been adjusted:
    - Slide thumbnail resized to ~30% of notes page height
    - Notes text area expanded below
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    if len(prs.slides) == 0:
        print("CRITICAL: No slides in presentation")
        print("REWARD: 0.0")
        return 0.0

    # We check all slides that have notes slides; the task says "notes page layout"
    # which typically means all notes pages should be adjusted.
    # We verify across all slides and take the average, but require consistency.

    slide_image_heights = []
    notes_body_tops = []
    notes_body_heights = []

    for i, slide in enumerate(prs.slides):
        try:
            ns = slide.notes_slide
            slide_img_ph = None
            body_ph = None
            for ph in ns.placeholders:
                ph_type_str = str(ph.placeholder_format.type)
                # Match by type: SLIDE_IMAGE (101) or BODY (2)
                if 'SLIDE_IMAGE' in ph_type_str:
                    slide_img_ph = ph
                elif 'BODY' in ph_type_str:
                    body_ph = ph

            if slide_img_ph is not None:
                slide_image_heights.append(slide_img_ph.height)
            if body_ph is not None:
                notes_body_tops.append(body_ph.top)
                notes_body_heights.append(body_ph.height)
        except Exception as e:
            print(f"WARNING: Could not read notes for slide {i+1}: {e}")

    if not slide_image_heights:
        print("CRITICAL: No slide image placeholders found in any notes slide")
        print("REWARD: 0.0")
        return 0.0

    # Use the first slide's values as representative (they should all be consistent)
    img_h = slide_image_heights[0]
    img_ratio = img_h / NOTES_PAGE_HEIGHT

    print(f"INFO: Slide image height = {img_h} EMU, ratio = {img_ratio:.4f} ({img_ratio*100:.1f}%)")
    print(f"INFO: Initial slide image height was {INITIAL_SLIDE_IMAGE_HEIGHT} EMU ({INITIAL_SLIDE_IMAGE_HEIGHT/NOTES_PAGE_HEIGHT*100:.1f}%)")

    # Component 1: Slide image height is ~30% of page (0.4 points)
    # Must be REDUCED from initial (~37.5%) toward ~30%.
    # This is the primary task change.
    try:
        # Use tolerance of 1% for "not reduced" check (handles minor EMU rounding)
        reduction_threshold = INITIAL_SLIDE_IMAGE_HEIGHT * 0.99
        if img_h >= reduction_threshold:
            # Not meaningfully reduced — no credit
            print(f"FAIL: Component 1 — Slide image not reduced (height={img_h}, threshold={reduction_threshold})")
        elif FULL_CREDIT_LOW <= img_ratio <= FULL_CREDIT_HIGH:
            print(f"PASS: Component 1 — Slide image height is {img_ratio*100:.1f}% of page (within 28-32% target) (0.4 pts)")
            total_score += 0.4
        elif PARTIAL_CREDIT_LOW <= img_ratio < FULL_CREDIT_LOW:
            partial = 0.2
            print(f"PARTIAL: Component 1 — Slide image height is {img_ratio*100:.1f}% (below target range, partial credit) ({partial} pts)")
            total_score += partial
        elif FULL_CREDIT_HIGH < img_ratio <= PARTIAL_CREDIT_HIGH:
            partial = 0.2
            print(f"PARTIAL: Component 1 — Slide image height is {img_ratio*100:.1f}% (above target range, partial credit) ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — Slide image height ratio {img_ratio*100:.1f}% is outside acceptable range (20-40%)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Notes body top position moved up (0.3 points)
    # The notes body should start closer to the top since the slide image is smaller.
    # Initial: 4343400, Golden: 3657600
    try:
        if notes_body_tops:
            body_top = notes_body_tops[0]
            print(f"INFO: Notes body top = {body_top} EMU, initial was {INITIAL_NOTES_BODY_TOP}")

            if body_top >= INITIAL_NOTES_BODY_TOP:
                print(f"FAIL: Component 2 — Notes body not moved up (top={body_top}, initial={INITIAL_NOTES_BODY_TOP})")
            else:
                # Must have moved up by a meaningful amount (at least 300000 EMU ~ 0.33 inch)
                moved_by = INITIAL_NOTES_BODY_TOP - body_top
                if moved_by >= 300000:
                    print(f"PASS: Component 2 — Notes body moved up by {moved_by} EMU (0.3 pts)")
                    total_score += 0.3
                else:
                    print(f"PARTIAL: Component 2 — Notes body moved up by only {moved_by} EMU (needs >= 300000)")
                    total_score += 0.1
        else:
            print(f"FAIL: Component 2 — No notes body placeholder found")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Notes body height expanded (0.3 points)
    # The notes body should be taller to fill the space freed by the smaller slide image.
    # Initial: 4114800, Golden: 4799013
    try:
        if notes_body_heights:
            body_h = notes_body_heights[0]
            print(f"INFO: Notes body height = {body_h} EMU, initial was {INITIAL_NOTES_BODY_HEIGHT}")

            if body_h <= INITIAL_NOTES_BODY_HEIGHT:
                print(f"FAIL: Component 3 — Notes body not expanded (height={body_h}, initial={INITIAL_NOTES_BODY_HEIGHT})")
            else:
                expanded_by = body_h - INITIAL_NOTES_BODY_HEIGHT
                if expanded_by >= 300000:
                    print(f"PASS: Component 3 — Notes body expanded by {expanded_by} EMU (0.3 pts)")
                    total_score += 0.3
                else:
                    print(f"PARTIAL: Component 3 — Notes body expanded by only {expanded_by} EMU (needs >= 300000)")
                    total_score += 0.1
        else:
            print(f"FAIL: Component 3 — No notes body placeholder found")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Check consistency across slides
    if len(slide_image_heights) > 1:
        unique_heights = set(slide_image_heights)
        if len(unique_heights) == 1:
            print(f"INFO: All {len(slide_image_heights)} slides have consistent slide image height")
        else:
            print(f"WARNING: Inconsistent slide image heights across slides: {unique_heights}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(FILE_PATH)
