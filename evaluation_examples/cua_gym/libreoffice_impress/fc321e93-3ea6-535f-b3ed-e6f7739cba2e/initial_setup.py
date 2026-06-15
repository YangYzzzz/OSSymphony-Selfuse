"""
Initial Setup: Lecture presentation with default notes page layout
Task ID: impress_ndo_026
Domain: libreoffice_impress

Creates a Lecture.pptx presentation with 5 slides containing realistic lecture
content. The notes page layout uses the default proportions (~50% slide thumbnail).
Opens the file in LibreOffice Impress.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_026'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Introduction to Machine Learning"
    slide1.placeholders[1].text = "CS 4780 — Fall 2025\nProfessor Elena Vasquez"

    notes1 = slide1.notes_slide.notes_text_frame
    notes1.text = (
        "Welcome students to the first lecture of CS 4780. "
        "Today we will cover the fundamental concepts of machine learning, "
        "including supervised and unsupervised learning paradigms. "
        "Remind students about office hours: Tuesdays 2-4pm in Gates Hall 459."
    )

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Course Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Supervised Learning"
    p2a = body2.add_paragraph()
    p2a.text = "Unsupervised Learning"
    p2a.level = 0
    p2b = body2.add_paragraph()
    p2b.text = "Reinforcement Learning"
    p2b.level = 0
    p2c = body2.add_paragraph()
    p2c.text = "Neural Networks & Deep Learning"
    p2c.level = 0
    p2d = body2.add_paragraph()
    p2d.text = "Practical Applications & Ethics"
    p2d.level = 0

    notes2 = slide2.notes_slide.notes_text_frame
    notes2.text = (
        "Spend approximately 5 minutes on this slide. "
        "Emphasize that the course follows a progression from classical methods "
        "to modern deep learning. Mention that the midterm will cover supervised "
        "and unsupervised learning, while the final includes all topics. "
        "Reference textbook chapters 1-3 for background reading."
    )

    # --- Slide 3: Supervised Learning ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Supervised Learning"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Training data consists of input-output pairs"
    p3a = body3.add_paragraph()
    p3a.text = "Classification: predicting discrete labels"
    p3a.level = 1
    p3b = body3.add_paragraph()
    p3b.text = "Regression: predicting continuous values"
    p3b.level = 1
    p3c = body3.add_paragraph()
    p3c.text = "Key algorithms: Linear Regression, SVM, Decision Trees"
    p3c.level = 0
    p3d = body3.add_paragraph()
    p3d.text = "Evaluation metrics: accuracy, precision, recall, F1-score"
    p3d.level = 0

    notes3 = slide3.notes_slide.notes_text_frame
    notes3.text = (
        "This is the core conceptual slide. Use the whiteboard to draw the "
        "decision boundary example for binary classification. "
        "Ask students: 'What happens when classes overlap?' "
        "Transition to the spam detection example — show how features like "
        "word frequency map to binary labels. "
        "Estimated time: 15 minutes including questions."
    )

    # --- Slide 4: Unsupervised Learning ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Unsupervised Learning"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "No labeled outputs — discover hidden structure"
    p4a = body4.add_paragraph()
    p4a.text = "Clustering: K-Means, DBSCAN, Hierarchical"
    p4a.level = 1
    p4b = body4.add_paragraph()
    p4b.text = "Dimensionality Reduction: PCA, t-SNE, UMAP"
    p4b.level = 1
    p4c = body4.add_paragraph()
    p4c.text = "Applications: customer segmentation, anomaly detection"
    p4c.level = 0

    notes4 = slide4.notes_slide.notes_text_frame
    notes4.text = (
        "Contrast with supervised learning — no ground truth labels. "
        "Demo the K-Means visualization from the course website. "
        "Mention that PCA will be covered in depth in Week 6. "
        "Highlight real-world use case: Netflix recommendation clusters. "
        "Assignment 1 will involve implementing K-Means from scratch."
    )

    # --- Slide 5: Next Steps ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Next Steps & Assignments"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Reading: Bishop Ch. 1-2, Murphy Ch. 1"
    p5a = body5.add_paragraph()
    p5a.text = "Assignment 1: Due September 15, 2025"
    p5a.level = 0
    p5b = body5.add_paragraph()
    p5b.text = "Lab Session: Thursday 3-5pm, Phillips 318"
    p5b.level = 0
    p5c = body5.add_paragraph()
    p5c.text = "Office Hours: Tuesday 2-4pm, Gates 459"
    p5c.level = 0
    p5d = body5.add_paragraph()
    p5d.text = "Discussion Forum: Piazza (access code: ml4780f25)"
    p5d.level = 0

    notes5 = slide5.notes_slide.notes_text_frame
    notes5.text = (
        "Wrap up by emphasizing the reading assignment. "
        "Bishop chapters provide mathematical foundations while Murphy "
        "offers a more practical perspective. "
        "Remind students that Assignment 1 requires Python and NumPy — "
        "direct those who need help to the TA prep session on Friday. "
        "End with: 'Questions? See you Thursday for the lab.'"
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
