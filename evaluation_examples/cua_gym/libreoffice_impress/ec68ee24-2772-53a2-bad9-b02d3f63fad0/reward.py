"""
Reward Script: Python Workshop slides with hands-on exercise sections
Task ID: impress_wf_021
Domain: libreoffice_impress
Scoring:
  C1: File on Desktop with 12 slides (0.15)
  C2: Slide 1 title 'Python for Beginners' (0.10)
  C3: Slide 2 agenda list (0.10)
  C4: Lesson slides code blocks in monospace on gray rounded rects (0.20)
  C5: Exercise slides yellow background #FFF9C4 (0.15)
  C6: Exercise slides 'EXERCISE' label in orange (0.10)
  C7: Lesson slides Fade transitions (0.10)
  C8: Exercise slides Push transitions (0.10)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_021'

def check_transition(pptx_path, slide_idx, expected_type):
    """slide_idx is 0-based. expected_type: 'fade', 'push', etc."""
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        try:
            with zf.open(f'ppt/slides/slide{slide_idx + 1}.xml') as f:
                root = ET.parse(f).getroot()
                tr = root.find('.//p:transition', ns)
                if tr is not None:
                    return tr.find(f'.//p:{expected_type}', ns) is not None
        except KeyError:
            pass
    return False


def get_slide_background_rgb(slide):
    """Return background RGB as string like 'FFF9C4', or None."""
    fill = slide.background.fill
    if fill.type is not None and fill.type == 1:  # solid fill
        return str(fill.fore_color.rgb)
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: File exists on Desktop with exactly 12 slides (0.15 points)
    try:
        num_slides = len(prs.slides)
        if num_slides == 12:
            print(f"PASS: Component 1 — 12 slides found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — expected 12 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 1 title contains 'Python for Beginners' (0.10 points)
    try:
        slide1 = prs.slides[0]
        all_text = ""
        for shape in slide1.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    all_text += para.text + " "
        if "python for beginners" in all_text.lower():
            print(f"PASS: Component 2 — 'Python for Beginners' found on slide 1 (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 — 'Python for Beginners' not found on slide 1. Text: {all_text[:100]}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 2 has agenda items (at least 4 distinct paragraphs with text) (0.10 points)
    try:
        slide2 = prs.slides[1]
        agenda_items = []
        for shape in slide2.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t and "agenda" not in t.lower():
                        agenda_items.append(t)
        if len(agenda_items) >= 4:
            print(f"PASS: Component 3 — Slide 2 agenda has {len(agenda_items)} items (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 3 — Slide 2 agenda has only {len(agenda_items)} items, expected >=4")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Lesson slides (3,4,6,7,9,10) have code in monospace font inside
    # light gray (#F5F5F5) rounded rectangle shapes (0.20 points)
    # Check: at least 4 of 6 lesson slides have a rounded rect with F5F5F5 fill
    # AND monospace font (Courier New, Consolas, etc.) inside
    try:
        lesson_indices = [2, 3, 5, 6, 8, 9]  # 0-based
        lesson_pass_count = 0
        for si in lesson_indices:
            if si >= len(prs.slides):
                continue
            slide = prs.slides[si]
            found_code_block = False
            for shape in slide.shapes:
                # Check for auto shape (rounded rectangle) with gray fill
                has_gray_fill = False
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or shape.shape_type == 1:
                        sf = shape.fill
                        if sf.type is not None and sf.type == 1:
                            clr = str(sf.fore_color.rgb)
                            if clr in ('F5F5F5', 'EFEFEF', 'F0F0F0', 'E8E8E8', 'EEEEEE', 'F2F2F2', 'E5E5E5'):
                                has_gray_fill = True
                except:
                    pass

                if has_gray_fill and shape.has_text_frame:
                    # Check for monospace font
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            fname = (run.font.name or "").lower()
                            if any(m in fname for m in ['courier', 'mono', 'consola', 'menlo', 'fira code']):
                                found_code_block = True
                                break
                        if found_code_block:
                            break
                if found_code_block:
                    break
            if found_code_block:
                lesson_pass_count += 1

        # Award points proportionally: need at least 4 of 6
        if lesson_pass_count >= 4:
            print(f"PASS: Component 4 — {lesson_pass_count}/6 lesson slides have code blocks (0.20 pts)")
            total_score += 0.20
        elif lesson_pass_count >= 2:
            partial = round(0.10 * lesson_pass_count / 6, 2)
            print(f"PARTIAL: Component 4 — {lesson_pass_count}/6 lesson slides have code blocks ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — only {lesson_pass_count}/6 lesson slides have code blocks")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Exercise slides (5, 8, 11 = indices 4, 7, 10) have yellow #FFF9C4 background (0.15 points)
    try:
        exercise_indices = [4, 7, 10]  # 0-based
        yellow_count = 0
        for si in exercise_indices:
            if si >= len(prs.slides):
                continue
            bg = get_slide_background_rgb(prs.slides[si])
            if bg is not None and bg.upper() == 'FFF9C4':
                yellow_count += 1
            else:
                print(f"  Slide {si+1} background: {bg}")

        if yellow_count == 3:
            print(f"PASS: Component 5 — all 3 exercise slides have #FFF9C4 background (0.15 pts)")
            total_score += 0.15
        elif yellow_count >= 1:
            partial = round(0.15 * yellow_count / 3, 2)
            print(f"PARTIAL: Component 5 — {yellow_count}/3 exercise slides have yellow bg ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 — no exercise slides have #FFF9C4 background")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Exercise slides have 'EXERCISE' label in orange (0.10 points)
    try:
        exercise_indices = [4, 7, 10]
        exercise_label_count = 0
        for si in exercise_indices:
            if si >= len(prs.slides):
                continue
            slide = prs.slides[si]
            found_label = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if 'exercise' in para.text.strip().lower():
                            for run in para.runs:
                                if 'exercise' in (run.text or '').strip().lower():
                                    try:
                                        if run.font.color.type is not None:
                                            rgb = str(run.font.color.rgb).upper()
                                            # Accept orange shades: FF9800, FF8C00, FFA500, etc.
                                            r = int(rgb[0:2], 16)
                                            g = int(rgb[2:4], 16)
                                            b = int(rgb[4:6], 16)
                                            if r >= 200 and g >= 100 and g <= 200 and b <= 50:
                                                found_label = True
                                    except:
                                        pass
                                if found_label:
                                    break
                        if found_label:
                            break
                if found_label:
                    break
            if found_label:
                exercise_label_count += 1

        if exercise_label_count == 3:
            print(f"PASS: Component 6 — all 3 exercise slides have orange 'EXERCISE' label (0.10 pts)")
            total_score += 0.10
        elif exercise_label_count >= 1:
            partial = round(0.10 * exercise_label_count / 3, 2)
            print(f"PARTIAL: Component 6 — {exercise_label_count}/3 have orange 'EXERCISE' ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 6 — no exercise slides have orange 'EXERCISE' label")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Lesson slides (3,4,6,7,9,10) have Fade transitions (0.10 points)
    try:
        lesson_slides = [2, 3, 5, 6, 8, 9]  # 0-based
        fade_count = 0
        for si in lesson_slides:
            if check_transition(file_path, si, 'fade'):
                fade_count += 1

        if fade_count >= 4:
            print(f"PASS: Component 7 — {fade_count}/6 lesson slides have Fade transition (0.10 pts)")
            total_score += 0.10
        elif fade_count >= 2:
            partial = round(0.05 * fade_count / 6, 2)
            print(f"PARTIAL: Component 7 — {fade_count}/6 lesson slides have Fade ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 7 — only {fade_count}/6 lesson slides have Fade transition")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    # Component 8: Exercise slides (5, 8, 11) have Push transitions (0.10 points)
    try:
        exercise_slides = [4, 7, 10]  # 0-based
        push_count = 0
        for si in exercise_slides:
            if check_transition(file_path, si, 'push'):
                push_count += 1

        if push_count >= 2:
            print(f"PASS: Component 8 — {push_count}/3 exercise slides have Push transition (0.10 pts)")
            total_score += 0.10
        elif push_count >= 1:
            print(f"PARTIAL: Component 8 — {push_count}/3 exercise slides have Push (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 8 — no exercise slides have Push transition")
    except Exception as e:
        print(f"ERROR: Component 8 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point: check Desktop path first, then home dir
persist_app_state()

file_path = f'{WORKDIR}/Desktop/Python_Workshop.pptx'
if not os.path.exists(file_path):
    # Fallback: check home dir with task_id name
    alt_path = f'{WORKDIR}/Desktop/python_workshop.pptx'
    if os.path.exists(alt_path):
        file_path = alt_path
    else:
        print(f"File not found: {file_path}")
        print("REWARD: 0.0")
        exit()

verify_task(file_path)
