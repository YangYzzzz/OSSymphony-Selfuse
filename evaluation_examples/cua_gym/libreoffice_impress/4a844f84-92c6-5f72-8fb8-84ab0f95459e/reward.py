"""
Reward Script: Remove all border lines from table on slide 4
Task ID: impress_tct_034
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): Majority of borders (>50%) are no-fill/zero-width
  Component 2 (0.5): ALL borders (100%) are no-fill/zero-width
"""

import os
import zipfile
import xml.etree.ElementTree as ET
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_034'

NS = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
BORDER_TAGS = ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']


def persist_app_state(domain: str):
    """Save any unsaved LibreOffice state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def is_border_removed(ln_element):
    """
    Check if a border line element represents a removed/invisible border.
    A border is considered removed if:
      - It has noFill child, OR
      - Its width ('w' attribute) is '0' or '0'
    """
    if ln_element is None:
        # No border element at all - could mean inherited/default
        return False
    w = ln_element.get('w', None)
    no_fill = ln_element.find('a:noFill', NS)
    if no_fill is not None:
        return True
    if w is not None and int(w) == 0:
        return True
    return False


def count_borders(pptx_path, slide_idx=3):
    """
    Count total and removed borders for the table on the specified slide.
    slide_idx is 0-based (slide 4 = index 3).
    Returns (total_borders, removed_borders).
    """
    slide_num = slide_idx + 1
    total = 0
    removed = 0

    with zipfile.ZipFile(pptx_path, 'r') as zf:
        slide_path = f'ppt/slides/slide{slide_num}.xml'
        if slide_path not in zf.namelist():
            print(f"FAIL: slide{slide_num}.xml not found in archive")
            return (0, 0)

        with zf.open(slide_path) as f:
            root = ET.parse(f).getroot()
            tcs = root.findall('.//a:tc', NS)
            if not tcs:
                print(f"FAIL: No table cells found on slide {slide_num}")
                return (0, 0)

            for tc in tcs:
                tcPr = tc.find('a:tcPr', NS)
                if tcPr is None:
                    # No tcPr means no explicit border settings
                    total += 4  # 4 borders per cell
                    continue
                for tag in BORDER_TAGS:
                    total += 1
                    ln = tcPr.find(tag, NS)
                    if is_border_removed(ln):
                        removed += 1

    return (total, removed)


def has_table_on_slide(pptx_path, slide_idx=3):
    """Check that slide 4 (0-indexed 3) has a table shape."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(pptx_path)
        if slide_idx >= len(prs.slides):
            return False
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                return True
        return False
    except Exception:
        return False


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: file exists and has a table on slide 4
    if not os.path.exists(file_path):
        print(f"CRITICAL: File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    if not has_table_on_slide(file_path, slide_idx=3):
        print("CRITICAL: No table found on slide 4")
        print("REWARD: 0.0")
        return 0.0

    # Count borders
    try:
        total_borders, removed_borders = count_borders(file_path, slide_idx=3)
        print(f"INFO: Total borders: {total_borders}, Removed: {removed_borders}")
    except Exception as e:
        print(f"CRITICAL: Could not analyze borders: {e}")
        print("REWARD: 0.0")
        return 0.0

    if total_borders == 0:
        print("CRITICAL: No border elements found to evaluate")
        print("REWARD: 0.0")
        return 0.0

    removal_ratio = removed_borders / total_borders

    # Component 1: Majority of borders (>50%) are removed (0.5 points)
    # This tests partial completion - the agent removed most borders
    try:
        if removal_ratio > 0.5:
            print(f"PASS: Component 1 — {removal_ratio*100:.1f}% borders removed (>50%) (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 1 — Only {removal_ratio*100:.1f}% borders removed (need >50%)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: ALL borders are removed (100%) (0.5 points)
    # This tests full completion - every single border is removed
    try:
        if removed_borders == total_borders:
            print(f"PASS: Component 2 — All {total_borders} borders removed (100%) (0.5 pts)")
            total_score += 0.5
        else:
            remaining = total_borders - removed_borders
            print(f"FAIL: Component 2 — {remaining}/{total_borders} borders still visible")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
