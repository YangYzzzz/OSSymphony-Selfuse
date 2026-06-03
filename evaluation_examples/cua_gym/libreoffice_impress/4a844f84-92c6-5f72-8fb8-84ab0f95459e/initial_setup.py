"""
Initial Setup: Create a 5-slide presentation with a bordered table on slide 4
Task ID: impress_tct_034
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn, nsmap
from lxml import etree
import copy

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_034'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_cell_text(cell, text, font_name="Calibri", font_size=Pt(12),
                  bold=False, color=None, alignment=None):
    """Helper to set text in a table cell with formatting."""
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if alignment:
        p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def set_cell_borders(cell, color_hex="000000", width_emu=12700):
    """Set visible borders on all sides of a table cell via XML."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    # Define the four border elements
    borders = {
        'a:lnL': 'left',
        'a:lnR': 'right',
        'a:lnT': 'top',
        'a:lnB': 'bottom',
    }

    for border_tag in borders:
        # Remove existing border element if any
        existing = tcPr.findall(qn(border_tag))
        for e in existing:
            tcPr.remove(e)

        # Create new border element
        ln = tcPr.makeelement(qn(border_tag), {'w': str(width_emu), 'cmpd': 'sng'})
        solidFill = ln.makeelement(qn('a:solidFill'), {})
        srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': color_hex})
        solidFill.append(srgbClr)
        ln.append(solidFill)
        tcPr.append(ln)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Clean Design Quarterly Report"
    slide1.placeholders[1].text = "Prepared by the Strategy & Analytics Team\nQ4 2025 Review"

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Revenue growth exceeded projections by 12% in Q4 2025"
    p2a = body2.add_paragraph()
    p2a.text = "Customer acquisition costs decreased by 8% quarter-over-quarter"
    p2a.level = 0
    p2b = body2.add_paragraph()
    p2b.text = "Three new enterprise partnerships signed in November"
    p2b.level = 0
    p2c = body2.add_paragraph()
    p2c.text = "Employee satisfaction scores reached an all-time high of 4.7/5.0"
    p2c.level = 0

    # --- Slide 3: Market Trends ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Market Trends & Insights"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Digital transformation spending increased 23% across key verticals"
    p3a = body3.add_paragraph()
    p3a.text = "AI-driven analytics tools are reshaping competitive landscapes"
    p3a.level = 0
    p3b = body3.add_paragraph()
    p3b.text = "Sustainability metrics now influence 40% of enterprise purchasing decisions"
    p3b.level = 0
    p3c = body3.add_paragraph()
    p3c.text = "Remote collaboration platforms show 15% higher adoption in mid-market"
    p3c.level = 0

    # --- Slide 4: Performance Data Table (3 columns x 5 rows, WITH borders) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Add title textbox
    txBox = slide4.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Regional Performance Summary"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x6E)

    # Create 5-row x 3-column table
    rows, cols = 5, 3
    table_shape = slide4.shapes.add_table(
        rows, cols,
        Inches(1.5), Inches(1.5),
        Inches(9), Inches(4)
    )
    table = table_shape.table

    # Set column widths
    table.columns[0].width = Inches(3)
    table.columns[1].width = Inches(3)
    table.columns[2].width = Inches(3)

    # Header row
    headers = ["Region", "Revenue (USD)", "Growth (%)"]
    header_color = RGBColor(0xFF, 0xFF, 0xFF)
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        set_cell_text(cell, h, bold=True, color=header_color,
                      alignment=PP_ALIGN.CENTER, font_size=Pt(14))
        # Header background
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x2E, 0x4A, 0x6E)

    # Data rows
    data = [
        ["North America", "$2,450,000", "14.2%"],
        ["Europe & UK", "$1,830,000", "9.7%"],
        ["Asia-Pacific", "$1,120,000", "22.1%"],
        ["Latin America", "$680,000", "11.5%"],
    ]
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            align = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            set_cell_text(cell, val, font_size=Pt(12), alignment=align)
            # Alternate row shading
            if r % 2 == 0:
                fill = cell.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF4)

    # Apply visible borders to ALL cells
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            set_cell_borders(cell, color_hex="2E4A6E", width_emu=12700)

    # --- Slide 5: Next Steps ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Next Steps & Action Items"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Finalize Q1 2026 budget allocations by January 15th"
    p5a = body5.add_paragraph()
    p5a.text = "Launch pilot program for AI-assisted customer support in APAC"
    p5a.level = 0
    p5b = body5.add_paragraph()
    p5b.text = "Schedule quarterly business review with new enterprise partners"
    p5b.level = 0
    p5c = body5.add_paragraph()
    p5c.text = "Present sustainability roadmap to the board of directors"
    p5c.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
