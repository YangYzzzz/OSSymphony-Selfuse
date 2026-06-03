"""
Reward Script: Restaurant menu presentation for digital display
Task ID: impress_wf_061
Domain: libreoffice_impress
Scoring:
  C1 (0.15) - File on Desktop with 8 slides
  C2 (0.15) - Cream (#FFF8E1) background on all slides
  C3 (0.15) - 'La Bella Cucina' on slide 1
  C4 (0.15) - Accent colors #2E7D32 and #C62828 used
  C5 (0.15) - Star shape(s) on slide 7
  C6 (0.10) - 2-column layout on slide 6 (wines + cocktails)
  C7 (0.15) - Fade transitions with 8-second auto-advance on all slides
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_061'
FILE_PATH = os.path.join(WORKDIR, 'Desktop', 'Menu_Display.pptx')


def get_slide_background_rgb(slide):
    """Get background color of a slide, handling inherited fills."""
    fill = slide.background.fill
    if fill.type == 1:  # SOLID
        return str(fill.fore_color.rgb)
    elif fill.type == 5:  # inherited from master
        try:
            master_fill = slide.slide_layout.slide_master.background.fill
            if master_fill.type == 1:
                return str(master_fill.fore_color.rgb)
        except Exception:
            pass
    return None


def collect_all_text_colors(slide):
    """Collect all RGB font colors used in a slide."""
    colors = set()
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.color.type is not None:
                            colors.add(str(run.font.color.rgb))
                    except Exception:
                        pass
    return colors


def check_transitions(pptx_path, num_slides):
    """Check Fade transition with 8-second auto-advance on all slides."""
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    pass_count = 0
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        for i in range(1, num_slides + 1):
            fname = f'ppt/slides/slide{i}.xml'
            try:
                with zf.open(fname) as f:
                    root = ET.parse(f).getroot()
                    tr = root.find('.//p:transition', ns)
                    if tr is not None:
                        has_fade = tr.find(f'.//p:fade', ns) is not None
                        adv_tm = tr.get('advTm', '')
                        # Accept 8000ms (8 seconds)
                        has_auto = adv_tm == '8000'
                        if has_fade and has_auto:
                            pass_count += 1
                        else:
                            print(f"  Slide {i}: fade={has_fade}, advTm={adv_tm}")
                    else:
                        print(f"  Slide {i}: no transition element")
            except KeyError:
                print(f"  Slide {i}: XML not found")
    return pass_count


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: file must exist
    if not os.path.exists(file_path):
        print(f"CRITICAL: File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: File on Desktop with exactly 8 slides (0.15 points)
    try:
        if num_slides == 8:
            print(f"PASS: Component 1 -- 8 slides found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 -- expected 8 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Cream (#FFF8E1) background on all slides (0.15 points)
    try:
        cream_count = 0
        for i, slide in enumerate(prs.slides):
            bg_color = get_slide_background_rgb(slide)
            if bg_color == 'FFF8E1':
                cream_count += 1
            else:
                print(f"  Slide {i+1} background: {bg_color}")
        if num_slides > 0 and cream_count == num_slides:
            print(f"PASS: Component 2 -- all {num_slides} slides have #FFF8E1 background (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 -- {cream_count}/{num_slides} slides have #FFF8E1 background")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: 'La Bella Cucina' on slide 1 (0.15 points)
    try:
        if num_slides >= 1:
            slide1_text = ""
            for shape in prs.slides[0].shapes:
                if hasattr(shape, 'text'):
                    slide1_text += " " + shape.text
            if 'La Bella Cucina' in slide1_text:
                print(f"PASS: Component 3 -- 'La Bella Cucina' found on slide 1 (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 3 -- 'La Bella Cucina' not found on slide 1. Text: {slide1_text[:100]}")
        else:
            print(f"FAIL: Component 3 -- no slides")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Accent colors #2E7D32 (green) and #C62828 (red) used (0.15 points)
    try:
        all_colors = set()
        for slide in prs.slides:
            all_colors.update(collect_all_text_colors(slide))
        has_green = '2E7D32' in all_colors
        has_red = 'C62828' in all_colors
        if has_green and has_red:
            print(f"PASS: Component 4 -- both accent colors #2E7D32 and #C62828 found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 -- green=#2E7D32 present={has_green}, red=#C62828 present={has_red}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Star shape(s) on slide 7 (0.15 points)
    try:
        if num_slides >= 7:
            slide7 = prs.slides[6]
            star_count = 0
            for shape in slide7.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    try:
                        ast = shape.auto_shape_type
                        # Check for any star variant (STAR_5_POINT=92, etc.)
                        if ast is not None and 'STAR' in str(ast):
                            star_count += 1
                    except ValueError:
                        pass
            if star_count > 0:
                print(f"PASS: Component 5 -- {star_count} star shape(s) on slide 7 (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 5 -- no star shapes found on slide 7")
        else:
            print(f"FAIL: Component 5 -- fewer than 7 slides")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # Component 6: 2-column layout on slide 6 (wines + cocktails) (0.10 points)
    # Verify there are at least 2 content text boxes positioned side by side
    try:
        if num_slides >= 6:
            slide6 = prs.slides[5]
            # Collect text boxes (excluding titles/headers at top)
            content_boxes = []
            for shape in slide6.shapes:
                if shape.shape_type == 17:  # TEXT_BOX
                    text = shape.text.strip().lower()
                    # Look for content text boxes with menu items
                    if ('wine' in text or 'chianti' in text or 'barolo' in text or
                        'cocktail' in text or 'aperol' in text or 'negroni' in text or
                        '$' in text):
                        content_boxes.append(shape)

            if len(content_boxes) >= 2:
                # Verify they are positioned in two columns (different left positions)
                lefts = [s.left for s in content_boxes]
                unique_lefts = set(lefts)
                if len(unique_lefts) >= 2:
                    print(f"PASS: Component 6 -- 2-column layout on slide 6 ({len(content_boxes)} content boxes, {len(unique_lefts)} columns) (0.10 pts)")
                    total_score += 0.10
                else:
                    print(f"FAIL: Component 6 -- content boxes not in 2 columns (left positions: {lefts})")
            else:
                print(f"FAIL: Component 6 -- expected 2+ content text boxes, found {len(content_boxes)}")
        else:
            print(f"FAIL: Component 6 -- fewer than 6 slides")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    # Component 7: Fade transitions with 8-second auto-advance on all slides (0.15 points)
    try:
        if num_slides > 0:
            pass_count = check_transitions(file_path, num_slides)
            if pass_count == num_slides:
                print(f"PASS: Component 7 -- all {num_slides} slides have Fade transition + 8s auto-advance (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 7 -- {pass_count}/{num_slides} slides have correct transitions")
        else:
            print(f"FAIL: Component 7 -- no slides")
    except Exception as e:
        print(f"ERROR: Component 7 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(FILE_PATH)
