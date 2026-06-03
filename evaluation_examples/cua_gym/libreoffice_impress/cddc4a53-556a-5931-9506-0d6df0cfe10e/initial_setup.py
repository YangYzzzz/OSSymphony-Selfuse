"""
Initial Setup: 7-slide training workshop deck - content textboxes in Calibri 14pt on slides 3-5
Task ID: osworld_impress_global_font_change_008
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_global_font_change_008'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_content(slide, title_text, content_lines, has_second_content=False, second_content_lines=None):
    """Add a title placeholder text and content textbox(es) to a slide."""
    # Title placeholder (index 0)
    title_ph = slide.placeholders[0]
    title_ph.text = title_text
    for para in title_ph.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(24)
            run.font.bold = True

    # Content textbox 1
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(8.5) if not has_second_content else Inches(4.2)
    height = Inches(4.5)

    txBox1 = slide.shapes.add_textbox(left, top, width, height)
    tf1 = txBox1.text_frame
    tf1.word_wrap = True

    for i, line in enumerate(content_lines):
        if i == 0:
            para = tf1.paragraphs[0]
        else:
            para = tf1.add_paragraph()
        para.text = line
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(14)

    if has_second_content and second_content_lines:
        left2 = Inches(4.8)
        txBox2 = slide.shapes.add_textbox(left2, top, width, height)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, line in enumerate(second_content_lines):
            if i == 0:
                para = tf2.paragraphs[0]
            else:
                para = tf2.add_paragraph()
            para.text = line
            for run in para.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(14)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Employee Development Workshop"
    slide1.placeholders[1].text = "Building Skills for the Modern Workplace\nQ1 2025 Training Series"
    for para in slide1.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(32)
            run.font.bold = True
    for para in slide1.placeholders[1].text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(18)

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Workshop Agenda"
    for para in slide2.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(24)
            run.font.bold = True
    content2 = slide2.placeholders[1]
    content2.text = ""
    agenda_items = [
        "Introduction & Team Overview",
        "Module 1: Communication Skills",
        "Module 2: Project Management",
        "Module 3: Technical Proficiency",
        "Module 4: Leadership & Collaboration",
        "Q&A and Wrap-Up",
    ]
    tf2 = content2.text_frame
    for i, item in enumerate(agenda_items):
        if i == 0:
            para = tf2.paragraphs[0]
        else:
            para = tf2.add_paragraph()
        para.text = item
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(14)

    # --- Slide 3: Module 1 - Communication Skills ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Module 1: Communication Skills"
    for para in slide3.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(24)
            run.font.bold = True
    content3 = slide3.placeholders[1]
    tf3 = content3.text_frame
    tf3.clear()
    comm_lines = [
        "Active listening techniques in team meetings",
        "Structuring clear and concise written reports",
        "Giving constructive feedback to colleagues",
        "Adapting communication style to different audiences",
        "Facilitating effective brainstorming sessions",
    ]
    for i, line in enumerate(comm_lines):
        if i == 0:
            para = tf3.paragraphs[0]
        else:
            para = tf3.add_paragraph()
        para.text = line
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(14)

    # --- Slide 4: Module 2 - Project Management ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Module 2: Project Management"
    for para in slide4.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(24)
            run.font.bold = True
    content4 = slide4.placeholders[1]
    tf4 = content4.text_frame
    tf4.clear()
    pm_lines = [
        "Setting SMART goals and milestones",
        "Resource allocation and workload balancing",
        "Risk identification and mitigation strategies",
        "Agile vs. Waterfall: choosing the right approach",
        "Tracking progress with dashboards and KPIs",
    ]
    for i, line in enumerate(pm_lines):
        if i == 0:
            para = tf4.paragraphs[0]
        else:
            para = tf4.add_paragraph()
        para.text = line
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(14)

    # --- Slide 5: Module 3 - Technical Proficiency ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Module 3: Technical Proficiency"
    for para in slide5.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(24)
            run.font.bold = True
    content5 = slide5.placeholders[1]
    tf5 = content5.text_frame
    tf5.clear()
    tech_lines = [
        "Advanced spreadsheet functions and data analysis",
        "Introduction to automation and scripting basics",
        "Effective use of collaboration tools (Slack, Jira)",
        "Data visualization best practices",
        "Cybersecurity awareness and safe computing habits",
    ]
    for i, line in enumerate(tech_lines):
        if i == 0:
            para = tf5.paragraphs[0]
        else:
            para = tf5.add_paragraph()
        para.text = line
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(14)

    # --- Slide 6: Module 4 - Leadership & Collaboration ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Module 4: Leadership & Collaboration"
    for para in slide6.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(24)
            run.font.bold = True
    content6 = slide6.placeholders[1]
    tf6 = content6.text_frame
    tf6.clear()
    lead_lines = [
        "Developing emotional intelligence in the workplace",
        "Building trust within cross-functional teams",
        "Conflict resolution frameworks and techniques",
        "Mentorship and coaching fundamentals",
        "Leading by example: values-driven decision making",
    ]
    for i, line in enumerate(lead_lines):
        if i == 0:
            para = tf6.paragraphs[0]
        else:
            para = tf6.add_paragraph()
        para.text = line
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(14)

    # --- Slide 7: Q&A and Wrap-Up ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Q&A and Wrap-Up"
    for para in slide7.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(24)
            run.font.bold = True
    content7 = slide7.placeholders[1]
    tf7 = content7.text_frame
    tf7.clear()
    qa_lines = [
        "Open floor for questions and discussion",
        "Key takeaways and action items",
        "Next steps: applying skills back on the job",
        "Feedback survey: help us improve future workshops",
        "Thank you for your participation!",
    ]
    for i, line in enumerate(qa_lines):
        if i == 0:
            para = tf7.paragraphs[0]
        else:
            para = tf7.add_paragraph()
        para.text = line
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(14)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
