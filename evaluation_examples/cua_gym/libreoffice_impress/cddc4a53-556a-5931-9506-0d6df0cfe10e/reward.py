"""
Reward Script: Set font of all content textboxes on slides 3, 4, and 5 to Georgia at 16pt
Task ID: osworld_impress_global_font_change_008
Domain: libreoffice_impress
Scoring:
  Component 1: Content placeholder on slide 3 uses Georgia at 16pt (0.35 pts)
  Component 2: Content placeholder on slide 4 uses Georgia at 16pt (0.35 pts)
  Component 3: Content placeholder on slide 5 uses Georgia at 16pt (0.30 pts)
  Total: 1.0
"""

import os

from pptx import Presentation
from pptx.util import Pt

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_global_font_change_008'

TARGET_FONT_NAME = 'Georgia'
TARGET_FONT_SIZE_PT = 16.0
# Placeholder index 1 identifies the content/body placeholder (not the title)
CONTENT_PH_IDX = 1
# Target slides (0-indexed: slides 3, 4, 5 → indices 2, 3, 4)
TARGET_SLIDE_INDICES = [2, 3, 4]


def get_content_shapes_on_slide(slide):
    """Return shapes that are content placeholders (placeholder index == 1)."""
    content_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.placeholder_format is not None:
            if shape.placeholder_format.idx == CONTENT_PH_IDX:
                content_shapes.append(shape)
    return content_shapes


def check_slide_content_font(prs, slide_idx):
    """
    Check that all non-empty runs in content placeholder on a given slide
    use Georgia font at 16pt.
    Returns (all_correct, total_runs_checked, failing_details)
    """
    slide = prs.slides[slide_idx]
    content_shapes = get_content_shapes_on_slide(slide)

    if not content_shapes:
        return False, 0, [f"No content placeholder found on slide {slide_idx + 1}"]

    total_runs = 0
    failing_details = []

    for shape in content_shapes:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not (run.text or "").strip():
                    continue
                total_runs += 1
                font_name = run.font.name
                font_size = run.font.size
                font_size_pt = None if font_size is None else font_size / 12700.0

                name_ok = (font_name == TARGET_FONT_NAME)
                size_ok = (font_size_pt is not None and abs(font_size_pt - TARGET_FONT_SIZE_PT) < 0.1)

                if not name_ok or not size_ok:
                    failing_details.append(
                        f"run='{run.text[:30]}' font={font_name} size={font_size_pt}pt"
                    )

    all_correct = (total_runs > 0 and len(failing_details) == 0)
    return all_correct, total_runs, failing_details


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Sanity check: ensure file has at least 5 slides
    if len(prs.slides) < 5:
        print(f"CRITICAL: Expected at least 5 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Content placeholder on slide 3 uses Georgia at 16pt (0.35 pts)
    try:
        all_correct, run_count, failing = check_slide_content_font(prs, slide_idx=2)
        if all_correct and run_count > 0:
            print(f"PASS: Component 1 — Slide 3 content runs ({run_count}) all use Georgia at 16pt (0.35 pts)")
            total_score += 0.35
        else:
            if failing:
                print(f"FAIL: Component 1 — Slide 3 content runs NOT all Georgia 16pt. Failures: {failing[:3]}")
            else:
                print(f"FAIL: Component 1 — Slide 3: no runs found or no content placeholder")
    except Exception as e:
        print(f"ERROR: Component 1 (Slide 3) — {e}")

    # Component 2: Content placeholder on slide 4 uses Georgia at 16pt (0.35 pts)
    try:
        all_correct, run_count, failing = check_slide_content_font(prs, slide_idx=3)
        if all_correct and run_count > 0:
            print(f"PASS: Component 2 — Slide 4 content runs ({run_count}) all use Georgia at 16pt (0.35 pts)")
            total_score += 0.35
        else:
            if failing:
                print(f"FAIL: Component 2 — Slide 4 content runs NOT all Georgia 16pt. Failures: {failing[:3]}")
            else:
                print(f"FAIL: Component 2 — Slide 4: no runs found or no content placeholder")
    except Exception as e:
        print(f"ERROR: Component 2 (Slide 4) — {e}")

    # Component 3: Content placeholder on slide 5 uses Georgia at 16pt (0.30 pts)
    try:
        all_correct, run_count, failing = check_slide_content_font(prs, slide_idx=4)
        if all_correct and run_count > 0:
            print(f"PASS: Component 3 — Slide 5 content runs ({run_count}) all use Georgia at 16pt (0.30 pts)")
            total_score += 0.30
        else:
            if failing:
                print(f"FAIL: Component 3 — Slide 5 content runs NOT all Georgia 16pt. Failures: {failing[:3]}")
            else:
                print(f"FAIL: Component 3 — Slide 5: no runs found or no content placeholder")
    except Exception as e:
        print(f"ERROR: Component 3 (Slide 5) — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the VM env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
