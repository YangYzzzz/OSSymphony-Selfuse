"""
Initial Setup: Create New_Hire_Onboarding.pptx with 10 blank-layout slides
Task ID: impress_ps_038
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_038'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # Standard 16:9 slide dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Use blank layout (index 6 = Blank in default template)
    blank_layout = prs.slide_layouts[6]

    slide_content = [
        {
            "title": "Welcome to Meridian Technologies",
            "subtitle": "New Hire Onboarding Program 2025",
        },
        {
            "title": "Company Overview",
            "bullets": [
                "Founded in 2008 with headquarters in Austin, Texas",
                "Over 3,200 employees across 12 global offices",
                "Annual revenue of $1.8 billion in fiscal year 2024",
                "Industry leader in cloud infrastructure solutions",
            ],
        },
        {
            "title": "Our Mission & Values",
            "bullets": [
                "Mission: Empowering businesses through innovative technology",
                "Integrity in every interaction with clients and colleagues",
                "Continuous learning and professional development",
                "Collaboration across teams and departments",
            ],
        },
        {
            "title": "Your First Week Schedule",
            "bullets": [
                "Monday: Orientation and IT setup with Sarah Chen (HR)",
                "Tuesday: Department introduction with your team lead",
                "Wednesday: Product training sessions (9:00 AM - 3:00 PM)",
                "Thursday: Security compliance and policy review",
                "Friday: Team lunch and one-on-one with your manager",
            ],
        },
        {
            "title": "Benefits & Compensation",
            "bullets": [
                "Health insurance coverage begins on Day 1",
                "401(k) matching up to 6% after 90 days",
                "20 days paid time off plus 10 company holidays",
                "Annual learning stipend of $2,500",
                "Employee stock purchase plan (ESPP)",
            ],
        },
        {
            "title": "IT Systems & Access",
            "bullets": [
                "Email: Microsoft Outlook (credentials from IT on Day 1)",
                "Project management: Jira and Confluence",
                "Communication: Slack workspace #meridian-general",
                "VPN access required for remote work days",
                "Two-factor authentication mandatory for all systems",
            ],
        },
        {
            "title": "Department Structure",
            "bullets": [
                "Engineering: Led by VP Marcus Johnson (Building C, Floor 3)",
                "Product: Led by VP Diana Reyes (Building A, Floor 2)",
                "Sales & Marketing: Led by VP Tom Richardson (Building B)",
                "Operations: Led by VP Kenji Watanabe (Building A, Floor 4)",
            ],
        },
        {
            "title": "Performance & Growth",
            "bullets": [
                "Quarterly check-ins with your direct manager",
                "Annual performance review cycle in December",
                "Internal mobility program for role transitions",
                "Mentorship pairing within your first 30 days",
                "Leadership development track for senior roles",
            ],
        },
        {
            "title": "Office Policies & Culture",
            "bullets": [
                "Hybrid work: 3 days in-office, 2 days remote",
                "Core hours: 10:00 AM - 4:00 PM local time",
                "Casual dress code except for client-facing meetings",
                "Monthly all-hands meeting on the first Friday",
                "Employee resource groups open to all team members",
            ],
        },
        {
            "title": "Key Contacts & Resources",
            "bullets": [
                "HR Business Partner: Sarah Chen (sarah.chen@meridian.com)",
                "IT Help Desk: helpdesk@meridian.com or ext. 4500",
                "Facilities: facilities@meridian.com for badge and access",
                "Employee Handbook: Available on the company intranet",
                "Emergency Contact: Security desk at ext. 9911",
            ],
        },
    ]

    for i, content in enumerate(slide_content):
        slide = prs.slides.add_slide(blank_layout)

        # Add title text box at top
        title_left = Inches(0.8)
        title_top = Inches(0.5)
        title_width = Inches(11.5)
        title_height = Inches(1.0)
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content["title"]
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        if "subtitle" in content:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.0))
            stf = sub_box.text_frame
            stf.word_wrap = True
            sp = stf.paragraphs[0]
            sp.text = content["subtitle"]
            sp.alignment = PP_ALIGN.LEFT
            sr = sp.runs[0]
            sr.font.name = "Arial"
            sr.font.size = Pt(24)
            sr.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

        if "bullets" in content:
            body_left = Inches(0.8)
            body_top = Inches(1.8)
            body_width = Inches(11.5)
            body_height = Inches(5.0)
            body_box = slide.shapes.add_textbox(body_left, body_top, body_width, body_height)
            btf = body_box.text_frame
            btf.word_wrap = True
            for j, bullet in enumerate(content["bullets"]):
                if j == 0:
                    bp = btf.paragraphs[0]
                else:
                    bp = btf.add_paragraph()
                bp.text = bullet
                bp.space_after = Pt(8)
                br = bp.runs[0]
                br.font.name = "Arial"
                br.font.size = Pt(18)
                br.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
