"""
Reward Script: Custom two-column layout with blue sidebar on slides 2, 4, 6, 8
Task ID: impress_ps_038
Domain: libreoffice_impress
Scoring:
  Component 1 (0.50) - Blue rectangle exists on target slides (2,4,6,8) with correct color and ~40% width
  Component 2 (0.30) - Target slides have title text positioned in the left sidebar area
  Component 3 (0.20) - Content text exists in the right ~60% area on target slides
  Gate: Non-target slides must remain without blue rectangles (penalty if violated)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_038'

# Target slides that should have the custom layout (0-indexed: 1, 3, 5, 7)
TARGET_SLIDE_INDICES = [1, 3, 5, 7]  # slides 2, 4, 6, 8
NON_TARGET_SLIDE_INDICES = [0, 2, 4, 6, 8, 9]  # slides 1, 3, 5, 7, 9, 10


def has_blue_rectangle(slide, slide_width):
    """
    Check if the slide has a rectangle shape with #2196F3 blue fill
    covering approximately the left 40% of the slide.
    Returns True if found, False otherwise.
    """
    expected_width_ratio = 0.40
    tolerance = 0.08  # allow 32%-48% width ratio

    for shape in slide.shapes:
        # Check if it's an auto shape (rectangle)
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            try:
                fill = shape.fill
                if fill.type is not None and fill.type == 1:  # SOLID fill
                    color_rgb = str(fill.fore_color.rgb).upper()
                    # Check for #2196F3 blue (allow minor variations)
                    if color_rgb == '2196F3':
                        # Check width is approximately 40% of slide width
                        width_ratio = shape.width / slide_width
                        if abs(width_ratio - expected_width_ratio) <= tolerance:
                            # Check it starts at or near the left edge
                            if shape.left <= slide_width * 0.05:
                                # Check it spans full height (or close)
                                return True
            except Exception:
                continue
    return False


def get_blue_rect_right_edge(slide, slide_width):
    """Return the right edge position of the blue rectangle, or None."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                fill = shape.fill
                if fill.type is not None and fill.type == 1:
                    color_rgb = str(fill.fore_color.rgb).upper()
                    if color_rgb == '2196F3':
                        return shape.left + shape.width
            except Exception:
                continue
    return None


def has_title_in_sidebar(slide, slide_width):
    """
    Check if there's a text shape positioned in the left sidebar area
    (within the left ~40%) that contains title-like text.
    """
    sidebar_right = slide_width * 0.50  # generous: text box can extend up to 50%

    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if not text:
                continue
            # Check if the shape's left edge is within the sidebar region
            # and the shape is not too wide (not a full-width content box)
            if shape.left < sidebar_right and shape.width < slide_width * 0.55:
                # This looks like a title in the sidebar
                return True
    return False


def has_content_on_right(slide, slide_width):
    """
    Check if there's a text shape with content positioned in the right ~60% area.
    """
    right_area_start = slide_width * 0.35  # content starts somewhere around 35-50%

    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if not text:
                continue
            # Content text box: left edge is in the right portion, has substantial text
            if shape.left >= right_area_start and len(text) > 20:
                return True
    return False


def count_shapes(slide):
    """Count total shapes on a slide."""
    return len(list(slide.shapes))


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    if num_slides < 10:
        print(f"FAIL: Expected at least 10 slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    slide_width = prs.slide_width

    # Gate check: Non-target slides must not have blue rectangles added
    # This is a precondition gate (true in both initial and golden), so no points awarded
    # If violated, we apply a penalty to the final score
    penalty = 0.0
    try:
        for idx in NON_TARGET_SLIDE_INDICES:
            slide = prs.slides[idx]
            if has_blue_rectangle(slide, slide_width):
                penalty += 0.15
                print(f"GATE FAIL: Slide {idx+1} unexpectedly has blue rectangle (penalty applied)")
            else:
                print(f"GATE OK: Slide {idx+1} has no blue rectangle (correct)")
    except Exception as e:
        print(f"ERROR: Gate check - {e}")

    # Component 1: Blue rectangle (#2196F3) on target slides 2, 4, 6, 8 (0.50 points)
    # Each target slide contributes 0.125 points
    try:
        comp1_score = 0.0
        for idx in TARGET_SLIDE_INDICES:
            slide = prs.slides[idx]
            if has_blue_rectangle(slide, slide_width):
                comp1_score += 0.125
                print(f"PASS: Slide {idx+1} has blue (#2196F3) rectangle sidebar (~40% width)")
            else:
                print(f"FAIL: Slide {idx+1} missing blue rectangle sidebar")
        total_score += comp1_score
        print(f"Component 1 subtotal: {comp1_score}/0.50")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Title text positioned in left sidebar area on target slides (0.30 points)
    # Each target slide contributes 0.075 points
    # The title text box must be narrower than full-width (positioned in sidebar, not spanning whole slide)
    try:
        comp2_score = 0.0
        for idx in TARGET_SLIDE_INDICES:
            slide = prs.slides[idx]
            if has_title_in_sidebar(slide, slide_width):
                comp2_score += 0.075
                print(f"PASS: Slide {idx+1} has title text in sidebar area")
            else:
                print(f"FAIL: Slide {idx+1} missing title text in sidebar area")
        total_score += comp2_score
        print(f"Component 2 subtotal: {comp2_score}/0.30")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Content text in right ~60% area on target slides (0.20 points)
    # Each target slide contributes 0.05 points
    try:
        comp3_score = 0.0
        for idx in TARGET_SLIDE_INDICES:
            slide = prs.slides[idx]
            if has_content_on_right(slide, slide_width):
                comp3_score += 0.05
                print(f"PASS: Slide {idx+1} has content text in right area")
            else:
                print(f"FAIL: Slide {idx+1} missing content text in right area")
        total_score += comp3_score
        print(f"Component 3 subtotal: {comp3_score}/0.20")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Apply penalty for gate violations
    total_score = max(0.0, total_score - penalty)

    final_score = round(min(total_score, 1.0), 4)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
