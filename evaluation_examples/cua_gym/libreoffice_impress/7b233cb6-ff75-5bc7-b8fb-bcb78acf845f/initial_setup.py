"""
Initial Setup: Research presentation with images not at top, text not underlined
Task ID: osworld_impress_image_top_underline_text_008
Domain: libreoffice_impress
"""

import os
import io
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_image_top_underline_text_008'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def make_methodology_image():
    """Create a simple methodology diagram image (PNG bytes)."""
    width, height = 480, 320
    img = Image.new('RGB', (width, height), color=(230, 240, 255))
    draw = ImageDraw.Draw(img)

    # Draw boxes representing methodology steps
    boxes = [
        (20, 60, 130, 120, (70, 130, 200), "Data\nCollection"),
        (170, 60, 280, 120, (70, 130, 200), "Data\nProcessing"),
        (320, 60, 430, 120, (70, 130, 200), "Analysis"),
        (20, 180, 130, 240, (50, 160, 100), "Validation"),
        (170, 180, 280, 240, (50, 160, 100), "Results"),
        (320, 180, 430, 240, (50, 160, 100), "Reporting"),
    ]

    for x1, y1, x2, y2, color, label in boxes:
        draw.rectangle([x1, y1, x2, y2], fill=color, outline=(20, 60, 120), width=2)
        cx = (x1 + x2) // 2
        cy = (y1 + y2) // 2
        draw.text((cx - 25, cy - 10), label, fill=(255, 255, 255))

    # Arrows between steps
    for sx, ex in [(130, 170), (280, 320)]:
        draw.line([(sx, 90), (ex, 90)], fill=(20, 60, 120), width=3)
        draw.polygon([(ex, 85), (ex, 95), (ex + 10, 90)], fill=(20, 60, 120))

    for sx, ex in [(130, 170), (280, 320)]:
        draw.line([(sx, 210), (ex, 210)], fill=(20, 60, 120), width=3)
        draw.polygon([(ex, 205), (ex, 215), (ex + 10, 210)], fill=(20, 60, 120))

    # Title
    draw.text((160, 15), "Research Methodology", fill=(20, 60, 120))

    buf = io.BytesIO()
    img.save(buf, format='PNG')
    buf.seek(0)
    return buf


def make_results_chart_image():
    """Create a simple bar chart image for results (PNG bytes)."""
    width, height = 480, 320
    img = Image.new('RGB', (width, height), color=(255, 250, 240))
    draw = ImageDraw.Draw(img)

    # Draw a bar chart
    data = [('Q1', 45), ('Q2', 72), ('Q3', 58), ('Q4', 89), ('Q5', 63)]
    bar_w = 60
    bar_gap = 20
    base_y = 260
    scale = 2.2

    colors = [
        (200, 80, 80),
        (80, 160, 80),
        (80, 120, 200),
        (200, 140, 60),
        (140, 80, 200),
    ]

    for i, (label, val) in enumerate(data):
        x = 40 + i * (bar_w + bar_gap)
        bar_h = int(val * scale)
        draw.rectangle([x, base_y - bar_h, x + bar_w, base_y], fill=colors[i], outline=(50, 50, 50))
        draw.text((x + 15, base_y + 5), label, fill=(50, 50, 50))
        draw.text((x + 18, base_y - bar_h - 18), str(val), fill=(50, 50, 50))

    # Axes
    draw.line([(30, 30), (30, base_y)], fill=(50, 50, 50), width=2)
    draw.line([(30, base_y), (460, base_y)], fill=(50, 50, 50), width=2)

    # Title
    draw.text((150, 10), "Experimental Results", fill=(20, 60, 120))

    buf = io.BytesIO()
    img.save(buf, format='PNG')
    buf.seek(0)
    return buf


def create_initial():
    prs = Presentation()
    # Standard 16:9 dimensions
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # ----------------------------------------------------------------
    # Slide 1: Title Slide with body text (NOT underlined)
    # ----------------------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Advances in Neural Network Architectures"
    subtitle = slide1.placeholders[1]
    subtitle.text = "A Comprehensive Research Overview\nDepartment of Computer Science, 2025"
    for para in subtitle.text_frame.paragraphs:
        for run in para.runs:
            run.font.underline = False
            run.font.size = Pt(20)

    # Body textbox (not underlined)
    txb1 = slide1.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1.5))
    tf1 = txb1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "This presentation summarizes three years of research into transformer-based architectures, attention mechanisms, and their practical applications in natural language processing and computer vision."
    for run in p1.runs:
        run.font.size = Pt(16)
        run.font.underline = False
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ----------------------------------------------------------------
    # Slide 2: Introduction with content text (NOT underlined)
    # ----------------------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Introduction & Motivation"
    content2 = slide2.placeholders[1]
    content2.text_frame.clear()
    lines2 = [
        "Deep learning has transformed how machines perceive and process information.",
        "Traditional architectures struggle with long-range dependencies in sequential data.",
        "Attention mechanisms offer a powerful alternative to recurrence-based approaches.",
        "Our research investigates three novel attention variants across benchmark datasets.",
        "We demonstrate consistent improvements over state-of-the-art baselines.",
    ]
    for i, line in enumerate(lines2):
        if i == 0:
            p = content2.text_frame.paragraphs[0]
        else:
            p = content2.text_frame.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.size = Pt(18)
            run.font.underline = False
            run.font.color.rgb = RGBColor(0x22, 0x22, 0x44)

    # ----------------------------------------------------------------
    # Slide 3: Methodology — with image NOT at top (placed in middle-left)
    # ----------------------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    # Title textbox
    txb3_title = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    txb3_title.text_frame.paragraphs[0].text = "Research Methodology"
    for run in txb3_title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

    # Image placed in the MIDDLE-LEFT area (NOT top-center) — agent needs to move it
    methodology_img = make_methodology_image()
    img_w = Inches(5.5)
    img_h = Inches(3.5)
    # Position: left=1inch, top=3.0inch (middle area, not top)
    pic3 = slide3.shapes.add_picture(methodology_img, Inches(1.0), Inches(3.0), img_w, img_h)

    # Caption text below image
    txb3_cap = slide3.shapes.add_textbox(Inches(1.0), Inches(6.6), Inches(11), Inches(0.7))
    txb3_cap.text_frame.paragraphs[0].text = "Figure 1: End-to-end research pipeline showing data flow from collection to reporting."
    for run in txb3_cap.text_frame.paragraphs[0].runs:
        run.font.size = Pt(13)
        run.font.italic = True
        run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # ----------------------------------------------------------------
    # Slide 4: Experimental Setup with content text (NOT underlined)
    # ----------------------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Experimental Setup"
    content4 = slide4.placeholders[1]
    content4.text_frame.clear()
    lines4 = [
        "Datasets: ImageNet-1K, COCO 2017, SQuAD v2, and WMT-2020 (EN-DE).",
        "Baseline models: ResNet-50, BERT-base, DeiT-S, and ViT-B/16.",
        "Training: 4x NVIDIA A100 GPUs, batch size 256, learning rate 1e-4.",
        "Evaluation metrics: Top-1 accuracy, F1 score, BLEU-4, and mAP@50.",
        "Ablation studies conducted for each architectural component.",
    ]
    for i, line in enumerate(lines4):
        if i == 0:
            p = content4.text_frame.paragraphs[0]
        else:
            p = content4.text_frame.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.size = Pt(18)
            run.font.underline = False
            run.font.color.rgb = RGBColor(0x22, 0x22, 0x44)

    # ----------------------------------------------------------------
    # Slide 5: Results — with chart image NOT at top (placed in center-right)
    # ----------------------------------------------------------------
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    # Title textbox
    txb5_title = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    txb5_title.text_frame.paragraphs[0].text = "Experimental Results"
    for run in txb5_title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

    # Image placed in CENTER-RIGHT area (NOT top) — agent needs to move it to top
    results_img = make_results_chart_image()
    img_w5 = Inches(5.5)
    img_h5 = Inches(3.5)
    # Position: left=6inch, top=2.8inch (not at top)
    pic5 = slide5.shapes.add_picture(results_img, Inches(6.0), Inches(2.8), img_w5, img_h5)

    # Caption
    txb5_cap = slide5.shapes.add_textbox(Inches(6.0), Inches(6.4), Inches(7), Inches(0.7))
    txb5_cap.text_frame.paragraphs[0].text = "Figure 2: Performance across five benchmark tasks (Q1-Q5)."
    for run in txb5_cap.text_frame.paragraphs[0].runs:
        run.font.size = Pt(13)
        run.font.italic = True
        run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # ----------------------------------------------------------------
    # Slide 6: Discussion with content text (NOT underlined)
    # ----------------------------------------------------------------
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Discussion & Analysis"
    content6 = slide6.placeholders[1]
    content6.text_frame.clear()
    lines6 = [
        "Multi-head attention variants consistently outperform single-head counterparts.",
        "Sparse attention reduces computational cost by 40% with less than 2% accuracy loss.",
        "Cross-modal attention enables strong performance on vision-language benchmarks.",
        "Limitations include higher memory requirements and complex training schedules.",
        "Future work will explore efficient fine-tuning strategies for resource-constrained devices.",
    ]
    for i, line in enumerate(lines6):
        if i == 0:
            p = content6.text_frame.paragraphs[0]
        else:
            p = content6.text_frame.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.size = Pt(18)
            run.font.underline = False
            run.font.color.rgb = RGBColor(0x22, 0x22, 0x44)

    # ----------------------------------------------------------------
    # Slide 7: Conclusion
    # ----------------------------------------------------------------
    slide7 = prs.slides.add_slide(prs.slide_layouts[0])
    slide7.shapes.title.text = "Conclusion"
    conclusion_ph = slide7.placeholders[1]
    conclusion_ph.text = "Thank you for your attention.\nQuestions & Discussion"
    for para in conclusion_ph.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(22)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
