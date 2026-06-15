"""
Reward Script: move methodology diagram (slide 3) to top-center, move results chart
(slide 5) to top, and underline all body/content text on slides 1, 2, 4, and 6.
Task ID: osworld_impress_image_top_underline_text_008
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 3 image moved to top-center  (0.25 pts)
  Component 2: Slide 5 image moved to top          (0.25 pts)
  Component 3: Body text on slides 1,2,4,6 underlined (0.50 pts, 0.125/slide)
Total: 1.0
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_image_top_underline_text_008'

# Tolerance for position comparisons (0.5%)
def is_approx_equal(a, b, tolerance=0.005):
    if a == b:
        return True
    if a == 0 or b == 0:
        return a == b
    return abs(a - b) / max(abs(a), abs(b)) <= tolerance


def get_all_text_runs(slide):
    """Return all non-empty runs from non-title text shapes on the slide."""
    runs = []
    for shape in slide.shapes:
        # Skip title placeholders (shape.name starts with 'Title')
        if 'Title' in shape.name and shape.shape_type == 14:
            continue
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or '').strip():
                        runs.append(run)
    return runs


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    slide_cx = slide_w // 2  # horizontal center of slide

    # --- Component 1: Slide 3 image moved to top-center (0.25 pts) ---
    # Golden: image center-x == slide_cx (centered), top < 1500000 EMU (~1.64 inches)
    # Initial: image left=914400, top=2743200 (not centered, middle of slide)
    try:
        slide3 = prs.slides[2]
        img_shape = None
        for shape in slide3.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_shape = shape
                break

        if img_shape is None:
            print("FAIL: Component 1 — No image found on slide 3")
        else:
            img_center_x = img_shape.left + img_shape.width // 2
            img_top = img_shape.top
            is_centered = is_approx_equal(img_center_x, slide_cx)
            is_at_top = img_top < 1500000  # less than ~1.64 inches from top

            if is_centered and is_at_top:
                print(f"PASS: Component 1 — Slide 3 image at top-center: "
                      f"center_x={img_center_x} (slide_cx={slide_cx}), top={img_top} (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 — Slide 3 image NOT at top-center: "
                      f"center_x={img_center_x} (expected ~{slide_cx}), "
                      f"top={img_top} (expected < 1500000), "
                      f"is_centered={is_centered}, is_at_top={is_at_top}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # --- Component 2: Slide 5 image moved to top (0.25 pts) ---
    # Golden: image top < 1500000 EMU and centered horizontally
    # Initial: image left=5486400, top=2560320 (right side, middle of slide)
    try:
        slide5 = prs.slides[4]
        img_shape5 = None
        for shape in slide5.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_shape5 = shape
                break

        if img_shape5 is None:
            print("FAIL: Component 2 — No image found on slide 5")
        else:
            img_center_x5 = img_shape5.left + img_shape5.width // 2
            img_top5 = img_shape5.top
            is_centered5 = is_approx_equal(img_center_x5, slide_cx)
            is_at_top5 = img_top5 < 1500000  # less than ~1.64 inches from top

            if is_at_top5 and is_centered5:
                print(f"PASS: Component 2 — Slide 5 image at top: "
                      f"center_x={img_center_x5}, top={img_top5} (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 2 — Slide 5 image NOT at top: "
                      f"center_x={img_center_x5} (expected ~{slide_cx}), "
                      f"top={img_top5} (expected < 1500000), "
                      f"is_centered={is_centered5}, is_at_top={is_at_top5}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # --- Component 3: Body text underlined on slides 1, 2, 4, 6 (0.50 pts) ---
    # 0.125 pts per slide; all non-title body/content runs must have underline=True
    slides_to_check = [(0, 1), (1, 2), (3, 4), (5, 6)]  # (0-indexed, 1-indexed)

    for slide_0idx, slide_1idx in slides_to_check:
        try:
            slide = prs.slides[slide_0idx]
            body_runs = get_all_text_runs(slide)

            if not body_runs:
                print(f"FAIL: Component 3 slide {slide_1idx} — No body runs found")
                continue

            all_underlined = all(run.font.underline is True for run in body_runs)
            num_underlined = sum(1 for run in body_runs if run.font.underline is True)

            if all_underlined:
                print(f"PASS: Component 3 slide {slide_1idx} — All {len(body_runs)} body runs underlined (0.125 pts)")
                total_score += 0.125
            else:
                print(f"FAIL: Component 3 slide {slide_1idx} — "
                      f"Only {num_underlined}/{len(body_runs)} runs underlined. "
                      f"Expected underline=True for all body runs.")
        except Exception as e:
            print(f"ERROR: Component 3 slide {slide_1idx} — {e}")

    final_score = round(min(total_score, 1.0), 4)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
