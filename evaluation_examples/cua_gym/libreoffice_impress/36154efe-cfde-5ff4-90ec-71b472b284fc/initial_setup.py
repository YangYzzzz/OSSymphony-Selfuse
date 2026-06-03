"""
Initial Setup: Create a large presentation with high-resolution images
Task ID: impress_fix_008
Domain: libreoffice_impress
"""

import os
import io
import shlex
import subprocess
import time
import random
import struct

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont, ImageFilter
import numpy as np

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_008'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

# --- Image generation helpers ---

GALLERY_THEMES = [
    ("Sunset Over Mountains", (255, 120, 50), (30, 30, 80)),
    ("Ocean Waves at Dawn", (60, 130, 200), (200, 180, 140)),
    ("Autumn Forest Path", (180, 100, 30), (40, 80, 30)),
    ("City Skyline at Night", (20, 20, 40), (255, 200, 80)),
    ("Lavender Fields", (140, 100, 180), (80, 160, 80)),
    ("Desert Sand Dunes", (220, 180, 120), (80, 140, 220)),
    ("Tropical Rainforest", (20, 80, 30), (60, 180, 80)),
    ("Snowy Mountain Peak", (220, 230, 240), (100, 130, 180)),
    ("Cherry Blossom Garden", (255, 180, 200), (60, 100, 50)),
    ("Northern Lights", (10, 20, 50), (50, 200, 120)),
    ("Coral Reef Underwater", (20, 60, 120), (255, 140, 60)),
    ("Bamboo Forest", (60, 120, 40), (180, 200, 140)),
    ("Volcanic Landscape", (60, 20, 10), (255, 80, 20)),
    ("Starry Night Sky", (10, 10, 30), (255, 255, 200)),
    ("Wildflower Meadow", (100, 180, 60), (255, 200, 60)),
    ("Frozen Lake", (160, 200, 230), (80, 100, 140)),
    ("Ancient Stone Bridge", (140, 130, 110), (60, 100, 60)),
    ("Waterfall in Jungle", (30, 80, 50), (180, 220, 255)),
    ("Rolling Green Hills", (60, 140, 40), (140, 200, 255)),
    ("Golden Wheat Field", (220, 190, 80), (80, 140, 220)),
    ("Misty Morning Lake", (160, 180, 190), (80, 120, 100)),
    ("Cliffside Coastline", (120, 100, 80), (50, 120, 180)),
    ("Redwood Forest", (100, 60, 30), (40, 80, 40)),
    ("Glacial Valley", (180, 210, 230), (80, 110, 80)),
    ("Sahara Oasis", (200, 180, 120), (40, 120, 60)),
    ("Alpine Meadow", (80, 160, 60), (200, 220, 255)),
    ("Monsoon Rainclouds", (60, 70, 80), (160, 180, 200)),
    ("Tulip Fields Netherlands", (255, 60, 60), (60, 160, 40)),
    ("Icy Fjord Norway", (100, 140, 180), (200, 220, 230)),
    ("Savanna at Twilight", (180, 120, 60), (40, 30, 60)),
    ("Pacific Island Beach", (40, 160, 200), (240, 220, 180)),
    ("Scottish Highlands", (80, 100, 60), (160, 170, 180)),
    ("Amazon River Basin", (40, 80, 40), (100, 160, 120)),
    ("Patagonian Steppe", (160, 140, 100), (120, 160, 200)),
    ("Japanese Zen Garden", (140, 140, 120), (80, 120, 80)),
    ("Mediterranean Coast", (40, 100, 180), (220, 200, 160)),
    ("Icelandic Geysers", (140, 160, 170), (200, 100, 40)),
    ("Australian Outback", (200, 120, 60), (60, 120, 200)),
    ("Canadian Rockies", (120, 140, 160), (200, 220, 240)),
    ("African Baobab Trees", (160, 140, 80), (60, 40, 30)),
]

def generate_high_res_image(title, color1, color2, width=3000, height=2250, seed=None):
    """Generate a high-resolution (300 DPI) photo-like image with noise for realistic file size."""
    if seed is not None:
        random.seed(seed)
        np.random.seed(seed % (2**31))

    # Create gradient as numpy array for speed
    ys = np.linspace(0, 1, height).reshape(-1, 1)
    c1 = np.array(color1, dtype=np.float64)
    c2 = np.array(color2, dtype=np.float64)
    gradient = (c1 * (1 - ys) + c2 * ys).astype(np.uint8)
    arr = np.broadcast_to(gradient[:, np.newaxis, :], (height, width, 3)).copy()

    # Add photographic noise (makes PNG much larger, simulating real photos)
    noise = np.random.randint(-25, 26, arr.shape, dtype=np.int16)
    arr = np.clip(arr.astype(np.int16) + noise, 0, 255).astype(np.uint8)

    img = Image.fromarray(arr)
    draw = ImageDraw.Draw(img)

    # Add geometric shapes for visual interest
    for _ in range(random.randint(20, 40)):
        x1 = random.randint(0, width)
        y1 = random.randint(0, height)
        size = random.randint(80, 500)
        opacity_color = (
            min(255, max(0, color1[0] + random.randint(-60, 60))),
            min(255, max(0, color1[1] + random.randint(-60, 60))),
            min(255, max(0, color1[2] + random.randint(-60, 60))),
        )
        shape_type = random.choice(['ellipse', 'rectangle', 'ellipse'])
        if shape_type == 'ellipse':
            draw.ellipse([x1, y1, x1 + size, y1 + int(size * 0.7)], fill=opacity_color)
        else:
            draw.rectangle([x1, y1, x1 + size, y1 + int(size * 0.6)], fill=opacity_color)

    # Add title text
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 72)
    except (IOError, OSError):
        font = ImageFont.load_default()

    bbox = draw.textbbox((0, 0), title, font=font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    tx = (width - tw) // 2
    ty = height - th - 100
    draw.text((tx + 3, ty + 3), title, fill=(0, 0, 0), font=font)
    draw.text((tx, ty), title, fill=(255, 255, 255), font=font)

    # Save as high-quality JPEG at 300 DPI (JPEG with noise = large file)
    buf = io.BytesIO()
    img.save(buf, format='JPEG', quality=98, dpi=(300, 300))
    buf.seek(0)
    return buf


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    random.seed(42)
    theme_idx = 0

    for slide_num in range(1, 26):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only layout

        # Set slide title
        if slide.shapes.title:
            slide.shapes.title.text = GALLERY_THEMES[theme_idx][0]
            for para in slide.shapes.title.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(28)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        # Determine number of images for this slide (1-3)
        if slide_num <= 8:
            num_images = 1
        elif slide_num <= 18:
            num_images = 2
        else:
            num_images = 3

        for img_idx in range(num_images):
            title, c1, c2 = GALLERY_THEMES[theme_idx]
            theme_idx = (theme_idx + 1) % len(GALLERY_THEMES)

            # Generate high-res image
            img_buf = generate_high_res_image(title, c1, c2,
                                              width=3000, height=2250,
                                              seed=slide_num * 100 + img_idx)

            # Position images based on count
            if num_images == 1:
                left = Inches(1.5)
                top = Inches(1.8)
                width = Inches(10.333)
                height = Inches(5.2)
            elif num_images == 2:
                left = Inches(0.5 + img_idx * 6.5)
                top = Inches(1.8)
                width = Inches(5.833)
                height = Inches(5.2)
            else:  # 3
                left = Inches(0.3 + img_idx * 4.3)
                top = Inches(1.8)
                width = Inches(3.933)
                height = Inches(5.2)

            pic = slide.shapes.add_picture(img_buf, left, top, width, height)

    prs.save(OUTPUT)
    file_size = os.path.getsize(OUTPUT)
    print(f'Initial file created: {OUTPUT} ({file_size / 1024 / 1024:.1f} MB)')
    print(f'Slides: {len(prs.slides)}, with high-resolution 300 DPI images')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
