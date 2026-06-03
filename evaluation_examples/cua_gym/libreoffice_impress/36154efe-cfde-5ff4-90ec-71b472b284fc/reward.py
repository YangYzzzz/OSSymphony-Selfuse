"""
Reward Script: Compress presentation images to 150 DPI
Task ID: impress_fix_008
Domain: libreoffice_impress
Scoring:
  Precondition gates: file exists, loadable, 25 slides, 49 images, positions preserved
  Component 1 (0.40): File size significantly reduced (under 30MB)
  Component 2 (0.35): All images compressed (blob sizes reduced below 1MB threshold)
  Component 3 (0.25): Image pixel dimensions reduced (from ~3000px to ~1500px width)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_008'

EXPECTED_SLIDES = 25
EXPECTED_IMAGES = 49
COMPRESSED_BLOB_THRESHOLD = 1_000_000  # 1 MB — originals are 4.4-5.4 MB each
FILE_SIZE_THRESHOLD = 30 * 1024 * 1024  # 30 MB per task requirement
# Original pixel width is 3000; compressed should be ~1500
ORIGINAL_MIN_PIXEL_WIDTH = 2000  # above this = uncompressed
COMPRESSED_MAX_PIXEL_WIDTH = 2000  # at or below this = compressed

# Expected positions (left, top) per slide pattern
SINGLE_IMG_POS = [(1371600, 1645920)]
DOUBLE_IMG_POS = [(457200, 1645920), (6400800, 1645920)]
TRIPLE_IMG_POS = [(274320, 1645920), (4206240, 1645920), (8138160, 1645920)]


def get_expected_positions(slide_idx):
    """Return expected (left, top) tuples for images on this slide (0-indexed)."""
    if slide_idx < 8:
        return SINGLE_IMG_POS
    elif slide_idx < 18:
        return DOUBLE_IMG_POS
    else:
        return TRIPLE_IMG_POS


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    file_size = os.path.getsize(file_path)

    # --- PRECONDITION GATES (no points, but must pass to proceed) ---

    # Gate 1: Correct slide count
    num_slides = len(prs.slides)
    if num_slides != EXPECTED_SLIDES:
        print(f"GATE FAIL: Expected {EXPECTED_SLIDES} slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    # Gate 2: Correct total image count
    total_images = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                total_images += 1
    if total_images != EXPECTED_IMAGES:
        print(f"GATE FAIL: Expected {EXPECTED_IMAGES} images, found {total_images}")
        print("REWARD: 0.0")
        return 0.0

    # Gate 3: Image positions preserved
    for slide_idx, slide in enumerate(prs.slides):
        expected_pos = get_expected_positions(slide_idx)
        actual_positions = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                actual_positions.append((shape.left, shape.top))
        expected_sorted = sorted(expected_pos)
        actual_sorted = sorted(actual_positions)
        if len(actual_sorted) != len(expected_sorted):
            print(f"GATE FAIL: Slide {slide_idx+1} image count mismatch")
            print("REWARD: 0.0")
            return 0.0
        for (el, et), (al, at) in zip(expected_sorted, actual_sorted):
            if el != 0 and abs(el - al) / max(abs(el), abs(al)) > 0.005:
                print(f"GATE FAIL: Slide {slide_idx+1} image position mismatch (left: {el} vs {al})")
                print("REWARD: 0.0")
                return 0.0
            if et != 0 and abs(et - at) / max(abs(et), abs(at)) > 0.005:
                print(f"GATE FAIL: Slide {slide_idx+1} image position mismatch (top: {et} vs {at})")
                print("REWARD: 0.0")
                return 0.0

    print("GATES PASSED: 25 slides, 49 images, all positions correct")

    # --- SCORING COMPONENTS (only changes between initial and golden) ---

    # Component 1: File size significantly reduced (0.40 points)
    # Initial is ~226 MB. Task requires under 30 MB.
    try:
        if file_size < FILE_SIZE_THRESHOLD:
            print(f"PASS: Component 1 — File size {file_size / 1024 / 1024:.1f} MB < 30 MB threshold (0.40 pts)")
            total_score += 0.40
        else:
            print(f"FAIL: Component 1 — File size {file_size / 1024 / 1024:.1f} MB >= 30 MB threshold")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: All images compressed - blob sizes reduced (0.35 points)
    # Original blobs are 4.4-5.4 MB each. Compressed should be well under 1 MB.
    # Award partial credit proportional to fraction of images compressed.
    try:
        compressed_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    blob_len = len(shape.image.blob)
                    if blob_len < COMPRESSED_BLOB_THRESHOLD:
                        compressed_count += 1

        if compressed_count == EXPECTED_IMAGES:
            print(f"PASS: Component 2 — All {EXPECTED_IMAGES} images have blob < 1 MB (0.35 pts)")
            total_score += 0.35
        elif compressed_count > 0:
            partial = 0.35 * (compressed_count / EXPECTED_IMAGES)
            print(f"PARTIAL: Component 2 — {compressed_count}/{EXPECTED_IMAGES} images compressed ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No images compressed (all blobs >= 1 MB)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Image pixel dimensions reduced (0.25 points)
    # Original images are 3000x2250 px. After 150 DPI compression they should be ~1500x1125.
    # Check that image pixel widths are below 2000 (clearly downscaled from 3000).
    try:
        import io
        from PIL import Image

        downscaled_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_blob = shape.image.blob
                    img = Image.open(io.BytesIO(img_blob))
                    w, h = img.size
                    if w <= COMPRESSED_MAX_PIXEL_WIDTH:
                        downscaled_count += 1

        if downscaled_count == EXPECTED_IMAGES:
            print(f"PASS: Component 3 — All {EXPECTED_IMAGES} images downscaled (width <= {COMPRESSED_MAX_PIXEL_WIDTH}px) (0.25 pts)")
            total_score += 0.25
        elif downscaled_count > 0:
            partial = 0.25 * (downscaled_count / EXPECTED_IMAGES)
            print(f"PARTIAL: Component 3 — {downscaled_count}/{EXPECTED_IMAGES} images downscaled ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No images downscaled (all widths > {COMPRESSED_MAX_PIXEL_WIDTH}px)")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
