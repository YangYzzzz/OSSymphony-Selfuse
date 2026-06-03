"""
Reward Script: Design Portfolio Presentation with Image Galleries
Task ID: impress_wf_024
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.15): File exists and has exactly 12 slides
  - Component 2 (0.10): Slide 1 title contains 'Design Portfolio' and 'Jane Smith'
  - Component 3 (0.10): Slide 2 has a circular-cropped image (ellipse geometry)
  - Component 4 (0.15): Slides 3-8 each have exactly 2 images
  - Component 5 (0.15): Slides 9-11 each have exactly 4 images (2x2 grid)
  - Component 6 (0.10): Images have gray (#E0E0E0) borders
  - Component 7 (0.15): Images have entrance animations (wipe from right)
  - Component 8 (0.10): Slide 12 has colored circle shapes with text initials
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_024'
FILE_PATH = os.path.join(WORKDIR, 'Desktop', 'Design_Portfolio.pptx')

# XML namespaces
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def count_pictures(slide):
    """Count picture shapes on a slide."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    return sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)


def check_image_borders(pptx_path, slide_indices):
    """Check if images on given slides have gray E0E0E0 borders via XML."""
    border_count = 0
    total_images = 0
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        for idx in slide_indices:
            fname = f'ppt/slides/slide{idx}.xml'
            try:
                with zf.open(fname) as f:
                    root = ET.parse(f).getroot()
                    # Find all pic elements
                    pics = root.findall(f'.//{{{NS_P}}}pic')
                    for pic in pics:
                        total_images += 1
                        ln = pic.find(f'.//{{{NS_A}}}ln')
                        if ln is not None:
                            srgb = ln.find(f'.//{{{NS_A}}}srgbClr')
                            if srgb is not None and srgb.get('val', '').upper() == 'E0E0E0':
                                border_count += 1
            except Exception:
                pass
    return border_count, total_images


def check_entrance_animations(pptx_path, slide_indices):
    """Check if slides have entrance animations (presetClass='entr') on images."""
    animated_slides = 0
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        for idx in slide_indices:
            fname = f'ppt/slides/slide{idx}.xml'
            try:
                with zf.open(fname) as f:
                    root = ET.parse(f).getroot()
                    timing = root.find(f'.//{{{NS_P}}}timing')
                    if timing is not None:
                        # Look for entrance animations (presetClass="entr")
                        ctns = timing.findall(f'.//{{{NS_P}}}cTn')
                        has_entrance = False
                        for ctn in ctns:
                            if ctn.get('presetClass') == 'entr':
                                has_entrance = True
                                break
                        if has_entrance:
                            animated_slides += 1
            except Exception:
                pass
    return animated_slides


def check_ellipse_image(pptx_path, slide_idx):
    """Check if a picture on given slide has ellipse preset geometry (circular crop)."""
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        fname = f'ppt/slides/slide{slide_idx}.xml'
        try:
            with zf.open(fname) as f:
                root = ET.parse(f).getroot()
                pics = root.findall(f'.//{{{NS_P}}}pic')
                for pic in pics:
                    prst_geom = pic.find(f'.//{{{NS_A}}}prstGeom')
                    if prst_geom is not None and prst_geom.get('prst') == 'ellipse':
                        return True
        except Exception:
            pass
    return False


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load presentation
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: File has exactly 12 slides (0.15 points)
    try:
        if num_slides == 12:
            print(f"PASS: Component 1 -- 12 slides found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 -- Expected 12 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Slide 1 title contains 'Design Portfolio' and 'Jane Smith' (0.10 points)
    try:
        if num_slides >= 1:
            slide1_text = ''
            for shape in prs.slides[0].shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        slide1_text += para.text + ' '
            slide1_text_lower = slide1_text.lower()
            if 'design portfolio' in slide1_text_lower and 'jane smith' in slide1_text_lower:
                print(f"PASS: Component 2 -- Slide 1 has 'Design Portfolio' and 'Jane Smith' (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 2 -- Slide 1 text: '{slide1_text.strip()[:100]}'")
        else:
            print("FAIL: Component 2 -- No slides available")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Slide 2 has circular-cropped image (ellipse geometry) (0.10 points)
    try:
        if num_slides >= 2:
            has_ellipse = check_ellipse_image(file_path, 2)
            if has_ellipse:
                print(f"PASS: Component 3 -- Slide 2 has ellipse-cropped image (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 3 -- Slide 2 has no ellipse-cropped image")
        else:
            print("FAIL: Component 3 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Slides 3-8 each have exactly 2 images (0.15 points)
    try:
        if num_slides >= 8:
            correct_count = 0
            for i in range(2, 8):  # slides 3-8 (0-indexed: 2-7)
                pic_count = count_pictures(prs.slides[i])
                if pic_count == 2:
                    correct_count += 1
                else:
                    print(f"  INFO: Slide {i+1} has {pic_count} images (expected 2)")
            # Award proportional credit
            proportion = correct_count / 6.0
            points = round(0.15 * proportion, 4)
            if correct_count == 6:
                print(f"PASS: Component 4 -- All slides 3-8 have 2 images each (0.15 pts)")
            else:
                print(f"PARTIAL: Component 4 -- {correct_count}/6 slides have 2 images ({points} pts)")
            total_score += points
        else:
            print(f"FAIL: Component 4 -- Not enough slides (need 8, have {num_slides})")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Slides 9-11 each have exactly 4 images (2x2 grid) (0.15 points)
    try:
        if num_slides >= 11:
            correct_count = 0
            for i in range(8, 11):  # slides 9-11 (0-indexed: 8-10)
                pic_count = count_pictures(prs.slides[i])
                if pic_count == 4:
                    correct_count += 1
                else:
                    print(f"  INFO: Slide {i+1} has {pic_count} images (expected 4)")
            proportion = correct_count / 3.0
            points = round(0.15 * proportion, 4)
            if correct_count == 3:
                print(f"PASS: Component 5 -- All slides 9-11 have 4 images each (0.15 pts)")
            else:
                print(f"PARTIAL: Component 5 -- {correct_count}/3 slides have 4 images ({points} pts)")
            total_score += points
        else:
            print(f"FAIL: Component 5 -- Not enough slides (need 11, have {num_slides})")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # Component 6: Images have gray (#E0E0E0) borders (0.10 points)
    # Check slides 3-11 (where content images are)
    try:
        check_slides = list(range(3, 12))  # slides 3-11
        border_count, total_images = check_image_borders(file_path, check_slides)
        if total_images > 0:
            proportion = border_count / total_images
            points = round(0.10 * proportion, 4)
            if proportion >= 0.8:
                print(f"PASS: Component 6 -- {border_count}/{total_images} images have E0E0E0 borders ({points} pts)")
            else:
                print(f"PARTIAL: Component 6 -- {border_count}/{total_images} images have E0E0E0 borders ({points} pts)")
            total_score += points
        else:
            print("FAIL: Component 6 -- No images found on slides 3-11")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    # Component 7: Images have entrance animations (0.15 points)
    # Check slides 3-11 (task says Fly In from Right on images)
    try:
        check_slides = list(range(3, 12))  # slides 3-11
        animated = check_entrance_animations(file_path, check_slides)
        expected = 9  # slides 3-11
        if animated > 0:
            proportion = min(animated / expected, 1.0)
            points = round(0.15 * proportion, 4)
            if animated >= expected:
                print(f"PASS: Component 7 -- {animated}/{expected} slides have entrance animations (0.15 pts)")
            else:
                print(f"PARTIAL: Component 7 -- {animated}/{expected} slides have entrance animations ({points} pts)")
            total_score += points
        else:
            print("FAIL: Component 7 -- No entrance animations found on slides 3-11")
    except Exception as e:
        print(f"ERROR: Component 7 -- {e}")

    # Component 8: Slide 12 has colored circle shapes with text initials (0.10 points)
    try:
        if num_slides >= 12:
            slide12 = prs.slides[11]
            oval_count = 0
            for shape in slide12.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    # Check if it's an oval with solid fill and has text
                    has_text = False
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                has_text = True
                                break
                    has_color = False
                    try:
                        if shape.fill.type == 1:  # SOLID
                            has_color = True
                    except Exception:
                        pass
                    if has_text and has_color:
                        oval_count += 1
            if oval_count >= 3:
                print(f"PASS: Component 8 -- Slide 12 has {oval_count} colored circles with text (0.10 pts)")
                total_score += 0.10
            elif oval_count > 0:
                points = round(0.10 * (oval_count / 3.0), 4)
                print(f"PARTIAL: Component 8 -- Slide 12 has {oval_count} colored circles ({points} pts)")
                total_score += points
            else:
                print("FAIL: Component 8 -- Slide 12 has no colored circles with text")
        else:
            print("FAIL: Component 8 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 8 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(FILE_PATH)
