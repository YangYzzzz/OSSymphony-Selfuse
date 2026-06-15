"""
Initial Setup: Presentation with audio that only plays on slide 1
Task ID: impress_fix_029
Domain: libreoffice_impress

Creates a 15-slide presentation with an embedded audio clip on slide 1
configured to stop at the end of slide 1 (NOT cross-slide playback).
"""

import os
import shlex
import shutil
import subprocess
import struct
import time
import wave
import zipfile
import copy
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_029'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
AUDIO_FILE = f'{WORKDIR}/background_music.mp3'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_dummy_mp3(path, duration_sec=10):
    """Create a minimal valid MP3 file (silent, very small)."""
    # Create a minimal MP3 frame (MPEG1 Layer3, 128kbps, 44100Hz, stereo)
    # Frame header: 0xFFFB9004 = sync + MPEG1, Layer3, 128kbps, 44100, stereo
    frame_header = bytes([0xFF, 0xFB, 0x90, 0x04])
    # Each frame at 128kbps/44100Hz is 417 bytes (or 418 with padding)
    frame_size = 417
    frame_data = frame_header + b'\x00' * (frame_size - 4)
    # ~38 frames per second at 44100Hz
    frames_per_sec = 38
    total_frames = frames_per_sec * duration_sec
    with open(path, 'wb') as f:
        # ID3v2 header (minimal)
        f.write(b'ID3')
        f.write(bytes([3, 0, 0]))  # version 2.3
        f.write(bytes([0, 0, 0, 0]))  # size = 0
        for _ in range(total_frames):
            f.write(frame_data)
    print(f"Created dummy MP3: {path} ({os.path.getsize(path)} bytes)")


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide themes for an "Ambiance Show"
    slide_data = [
        {"title": "Ambiance Show", "subtitle": "A Journey Through Light and Sound", "layout": 0,
         "bg": RGBColor(0x1A, 0x1A, 0x2E)},
        {"title": "Morning Mist", "content": "Soft light filters through the canopy\nDew drops catch the first rays of dawn\nBirds begin their morning chorus", "layout": 1,
         "bg": RGBColor(0x2C, 0x3E, 0x50)},
        {"title": "Ocean Breeze", "content": "Waves crash gently against weathered rocks\nSalt air carries the scent of seaweed\nGulls circle overhead in lazy spirals", "layout": 1,
         "bg": RGBColor(0x1B, 0x4F, 0x72)},
        {"title": "Desert Sunset", "content": "The horizon blazes with amber and gold\nSand dunes cast long purple shadows\nA lone coyote howls in the distance", "layout": 1,
         "bg": RGBColor(0x6C, 0x3A, 0x2A)},
        {"title": "Mountain Echo", "content": "Snow-capped peaks pierce the clouds\nEagles soar through crystal-clear air\nA stream cascades down granite walls", "layout": 1,
         "bg": RGBColor(0x34, 0x49, 0x5E)},
        {"title": "City Lights", "content": "Neon reflections dance in puddles\nTaxis weave through midnight traffic\nJazz drifts from an open window", "layout": 1,
         "bg": RGBColor(0x2D, 0x2D, 0x2D)},
        {"title": "Rainforest Canopy", "content": "Emerald leaves drip with moisture\nExotic flowers bloom in hidden clearings\nMonkeys chatter from branch to branch", "layout": 1,
         "bg": RGBColor(0x1D, 0x42, 0x2E)},
        {"title": "Arctic Aurora", "content": "Green curtains of light ripple across the sky\nIce crystals sparkle like scattered diamonds\nSilence stretches to the horizon", "layout": 1,
         "bg": RGBColor(0x0D, 0x1B, 0x2A)},
        {"title": "Autumn Woods", "content": "Crimson and gold leaves spiral downward\nMushrooms cluster at the base of old oaks\nA deer pauses at the edge of a clearing", "layout": 1,
         "bg": RGBColor(0x5D, 0x3A, 0x1A)},
        {"title": "Coral Reef", "content": "Tropical fish dart through anemone forests\nSea turtles glide with ancient grace\nSunlight filters through turquoise water", "layout": 1,
         "bg": RGBColor(0x00, 0x6B, 0x75)},
        {"title": "Lavender Fields", "content": "Purple rows stretch to the Provencal hills\nBees hum their industrious melodies\nA warm breeze carries floral perfume", "layout": 1,
         "bg": RGBColor(0x4A, 0x30, 0x6D)},
        {"title": "Thunderstorm", "content": "Lightning splits the charcoal sky\nRain hammers against windowpanes\nThunder rolls across the valley floor", "layout": 1,
         "bg": RGBColor(0x1C, 0x1C, 0x2A)},
        {"title": "Bamboo Grove", "content": "Tall stalks sway in synchronized rhythm\nDappled sunlight creates moving patterns\nA stone path leads to a hidden temple", "layout": 1,
         "bg": RGBColor(0x2E, 0x4A, 0x2E)},
        {"title": "Starry Night", "content": "The Milky Way arcs across the darkness\nMeteors streak brief trails of fire\nOwls call from the ancient pines", "layout": 1,
         "bg": RGBColor(0x0A, 0x0A, 0x23)},
        {"title": "Thank You", "subtitle": "Ambiance Show — Created with Care", "layout": 0,
         "bg": RGBColor(0x1A, 0x1A, 0x2E)},
    ]

    for i, sd in enumerate(slide_data):
        layout_idx = sd["layout"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Background
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = sd["bg"]

        # Title
        if slide.shapes.title:
            slide.shapes.title.text = sd["title"]
            for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.size = Pt(36)

        # Content or subtitle
        if "subtitle" in sd and len(slide.placeholders) > 1:
            slide.placeholders[1].text = sd["subtitle"]
            for run in slide.placeholders[1].text_frame.paragraphs[0].runs:
                run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
                run.font.size = Pt(20)
        elif "content" in sd and len(slide.placeholders) > 1:
            slide.placeholders[1].text = sd["content"]
            for para in slide.placeholders[1].text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
                    run.font.size = Pt(18)

    prs.save(OUTPUT)
    print(f"Presentation saved: {OUTPUT} with {len(prs.slides)} slides")
    return OUTPUT


def embed_audio_slide1_only(pptx_path, audio_path):
    """
    Embed audio on slide 1, configured to play only on that slide (not cross-slide).
    Uses direct XML manipulation to add audio with proper OOXML structure.
    """
    tmp_path = pptx_path + '.tmp'

    # Read the audio file bytes
    with open(audio_path, 'rb') as f:
        audio_bytes = f.read()

    with zipfile.ZipFile(pptx_path, 'r') as zin:
        with zipfile.ZipFile(tmp_path, 'w') as zout:
            # Copy all existing entries
            for item in zin.infolist():
                data = zin.read(item.filename)

                if item.filename == 'ppt/slides/slide1.xml':
                    # Modify slide1 to add audio shape
                    data = _add_audio_to_slide_xml(data, slide_only=True)
                elif item.filename == 'ppt/slides/_rels/slide1.xml.rels':
                    # Add relationship for audio
                    data = _add_audio_rel(data)
                elif item.filename == '[Content_Types].xml':
                    # Add mp3 content type
                    data = _add_mp3_content_type(data)

                zout.writestr(item, data)

            # Add the audio file to the package
            zout.writestr('ppt/media/background_music.mp3', audio_bytes)

    shutil.move(tmp_path, pptx_path)
    print(f"Audio embedded on slide 1 (single-slide playback)")


def _add_mp3_content_type(ct_data):
    """Add MP3 content type to [Content_Types].xml."""
    root = ET.fromstring(ct_data)
    ns = 'http://schemas.openxmlformats.org/package/2006/content-types'

    # Check if mp3 extension already exists
    for ext_el in root.findall(f'{{{ns}}}Default'):
        if ext_el.get('Extension') == 'mp3':
            return ct_data

    # Add mp3 default content type
    new_el = ET.SubElement(root, f'{{{ns}}}Default')
    new_el.set('Extension', 'mp3')
    new_el.set('ContentType', 'audio/mpeg')

    return ET.tostring(root, xml_declaration=True, encoding='UTF-8')


def _add_audio_rel(rel_data):
    """Add audio relationship to slide1 rels."""
    root = ET.fromstring(rel_data)
    ns = 'http://schemas.openxmlformats.org/package/2006/relationships'

    # Find next rId
    max_id = 0
    for rel in root.findall(f'{{{ns}}}Relationship'):
        rid = rel.get('Id', '')
        if rid.startswith('rId'):
            try:
                max_id = max(max_id, int(rid[3:]))
            except ValueError:
                pass

    audio_rid = f'rId{max_id + 1}'

    # Add relationship for audio
    new_rel = ET.SubElement(root, f'{{{ns}}}Relationship')
    new_rel.set('Id', audio_rid)
    new_rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/audio')
    new_rel.set('Target', '../media/background_music.mp3')

    # Store the rId for use in slide XML
    global AUDIO_RID
    AUDIO_RID = audio_rid

    return ET.tostring(root, xml_declaration=True, encoding='UTF-8')


AUDIO_RID = 'rId99'


def _add_audio_to_slide_xml(slide_data, slide_only=True):
    """Add audio shape element to slide XML.

    When slide_only=True: audio plays only on this slide (no cross-slide).
    When slide_only=False: audio plays across all slides.
    """
    namespaces = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main',
    }
    for prefix, uri in namespaces.items():
        ET.register_namespace(prefix, uri)
    # Also register common namespaces
    ET.register_namespace('mc', 'http://schemas.openxmlformats.org/markup-compatibility/2006')

    root = ET.fromstring(slide_data)

    # We'll use a simpler approach: add audio via the transition sound mechanism
    # For slide-only playback, we add a simple audio shape with timing limited to current slide

    # Find the spTree (shape tree)
    ns_p = namespaces['p']
    ns_a = namespaces['a']
    ns_r = namespaces['r']

    cSld = root.find(f'{{{ns_p}}}cSld')
    spTree = cSld.find(f'{{{ns_p}}}spTree')

    # Add a small audio icon shape (positioned in bottom-right corner)
    # We use pic element for the audio with special nvPicPr
    audio_shape_xml = f'''<p:sp xmlns:p="{ns_p}" xmlns:a="{ns_a}" xmlns:r="{ns_r}">
      <p:nvSpPr>
        <p:cNvPr id="9999" name="Audio: background_music">
          <a:hlinkClick r:id="" action="ppaction://media"/>
        </p:cNvPr>
        <p:cNvSpPr/>
        <p:nvPr>
          <a:audioFile r:link="{AUDIO_RID}"/>
        </p:nvPr>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm>
          <a:off x="8229600" y="5943600"/>
          <a:ext cx="457200" cy="457200"/>
        </a:xfrm>
        <a:prstGeom prst="rect">
          <a:avLst/>
        </a:prstGeom>
      </p:spPr>
      <p:txBody>
        <a:bodyPr/>
        <a:lstStyle/>
        <a:p>
          <a:r>
            <a:rPr lang="en-US" sz="800"/>
            <a:t>🔊</a:t>
          </a:r>
        </a:p>
      </p:txBody>
    </p:sp>'''

    audio_el = ET.fromstring(audio_shape_xml)
    spTree.append(audio_el)

    # Add timing for audio - play only on this slide (no cross-slide)
    # This is the key difference: numSld="0" means current slide only
    timing_xml = f'''<p:timing xmlns:p="{ns_p}" xmlns:a="{ns_a}">
      <p:tnLst>
        <p:par>
          <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
            <p:childTnLst>
              <p:seq concurrent="1" nextAc="seek">
                <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                  <p:childTnLst>
                    <p:par>
                      <p:cTn id="3" fill="hold">
                        <p:stCondLst>
                          <p:cond delay="0"/>
                        </p:stCondLst>
                        <p:childTnLst>
                          <p:par>
                            <p:cTn id="4" fill="hold">
                              <p:stCondLst>
                                <p:cond delay="0"/>
                              </p:stCondLst>
                              <p:childTnLst>
                                <p:audio>
                                  <p:cMediaNode numSld="1">
                                    <p:cTn id="5" fill="hold" display="0">
                                      <p:stCondLst>
                                        <p:cond delay="0"/>
                                      </p:stCondLst>
                                    </p:cTn>
                                    <p:tgtEl>
                                      <p:spTgt spid="9999"/>
                                    </p:tgtEl>
                                  </p:cMediaNode>
                                </p:audio>
                              </p:childTnLst>
                            </p:cTn>
                          </p:par>
                        </p:childTnLst>
                      </p:cTn>
                    </p:par>
                  </p:childTnLst>
                </p:cTn>
                <p:prevCondLst>
                  <p:cond evt="onPrev" delay="0">
                    <p:tgtEl>
                      <p:sldTgt/>
                    </p:tgtEl>
                  </p:cond>
                </p:prevCondLst>
                <p:nextCondLst>
                  <p:cond evt="onNext" delay="0">
                    <p:tgtEl>
                      <p:sldTgt/>
                    </p:tgtEl>
                  </p:cond>
                </p:nextCondLst>
              </p:seq>
            </p:childTnLst>
          </p:cTn>
        </p:par>
      </p:tnLst>
    </p:timing>'''

    timing_el = ET.fromstring(timing_xml)

    # Remove existing timing if present
    existing_timing = root.find(f'{{{ns_p}}}timing')
    if existing_timing is not None:
        root.remove(existing_timing)

    root.append(timing_el)

    return ET.tostring(root, xml_declaration=True, encoding='UTF-8')


def create_initial():
    # Step 1: Create dummy audio file
    create_dummy_mp3(AUDIO_FILE, duration_sec=10)

    # Step 2: Create presentation with 15 slides
    create_presentation()

    # Step 3: Embed audio on slide 1 with single-slide playback
    embed_audio_slide1_only(OUTPUT, AUDIO_FILE)

    # Step 4: Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
