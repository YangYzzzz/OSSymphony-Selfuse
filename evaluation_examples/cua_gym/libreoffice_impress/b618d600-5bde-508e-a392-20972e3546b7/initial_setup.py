"""
Initial Setup: Replace 'OldCorp' with 'NewCorp Technologies' across all slides
Task ID: impress_fix_092
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_092'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_shape(shape, text, font_size=Pt(14), bold=False, color=None, alignment=None):
    """Helper to set text on a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_paragraph(text_frame, text, font_size=Pt(14), bold=False, color=None, level=0):
    """Add a new paragraph to a text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    run = p.runs[0]
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return p


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide (has OldCorp in title) ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "OldCorp Annual Strategy Review"
    slide1.placeholders[1].text = "Fiscal Year 2025 | Confidential"

    # --- Slide 2: Company Overview (has OldCorp in body) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Company Overview"
    tf2 = slide2.placeholders[1].text_frame
    tf2.paragraphs[0].text = "OldCorp was founded in 2003 with a mission to transform enterprise logistics."
    add_paragraph(tf2, "Headquartered in San Francisco with offices in 12 countries.")
    add_paragraph(tf2, "Over 4,500 employees serving Fortune 500 clients worldwide.")

    # --- Slide 3: Mission & Values (has OldCorp in body) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Mission & Core Values"
    tf3 = slide3.placeholders[1].text_frame
    tf3.paragraphs[0].text = "At OldCorp, we believe in innovation through collaboration."
    add_paragraph(tf3, "Integrity: We hold ourselves to the highest ethical standards.")
    add_paragraph(tf3, "Excellence: Every product reflects our commitment to quality.")
    add_paragraph(tf3, "Sustainability: Building solutions that protect our planet.")

    # --- Slide 4: Financial Highlights ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Financial Highlights Q4 2025"
    tf4 = slide4.placeholders[1].text_frame
    tf4.paragraphs[0].text = "Revenue: $2.34B (+18% YoY)"
    add_paragraph(tf4, "EBITDA: $412M (17.6% margin)")
    add_paragraph(tf4, "Net Income: $289M (+22% YoY)")
    add_paragraph(tf4, "Free Cash Flow: $198M")

    # --- Slide 5: Market Position (has OldCorp in title) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "OldCorp Market Position"
    tf5 = slide5.placeholders[1].text_frame
    tf5.paragraphs[0].text = "Leader in enterprise logistics management (Gartner Magic Quadrant)."
    add_paragraph(tf5, "35% market share in North America.")
    add_paragraph(tf5, "Fastest-growing provider in APAC region (+42% YoY).")

    # --- Slide 6: Product Portfolio ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Product Portfolio"
    tf6 = slide6.placeholders[1].text_frame
    tf6.paragraphs[0].text = "LogiTrack Pro: End-to-end supply chain visibility"
    add_paragraph(tf6, "FleetCommand: Real-time fleet management and optimization")
    add_paragraph(tf6, "WarehouseIQ: AI-powered warehouse automation")
    add_paragraph(tf6, "RouteGenius: Dynamic routing and scheduling")

    # --- Slide 7: Customer Success (has OldCorp in body) ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Customer Success Stories"
    tf7 = slide7.placeholders[1].text_frame
    tf7.paragraphs[0].text = "Amazon reduced last-mile delivery costs by 23% using OldCorp solutions."
    add_paragraph(tf7, "Walmart improved warehouse throughput by 31% in 6 months.")
    add_paragraph(tf7, "DHL achieved 99.7% on-time delivery with FleetCommand.")

    # --- Slide 8: Technology Stack ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Technology Architecture"
    tf8 = slide8.placeholders[1].text_frame
    tf8.paragraphs[0].text = "Cloud-native microservices on AWS and Azure"
    add_paragraph(tf8, "ML pipeline: TensorFlow, PyTorch for demand forecasting")
    add_paragraph(tf8, "Real-time event streaming with Apache Kafka")
    add_paragraph(tf8, "GraphQL API layer for partner integrations")

    # --- Slide 9: Team & Leadership ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Executive Leadership"
    tf9 = slide9.placeholders[1].text_frame
    tf9.paragraphs[0].text = "CEO: Maria Chen (ex-Google, ex-McKinsey)"
    add_paragraph(tf9, "CTO: James Nakamura (15+ years in logistics tech)")
    add_paragraph(tf9, "CFO: Priya Sharma (ex-Goldman Sachs)")
    add_paragraph(tf9, "COO: David Okafor (ex-Amazon Logistics)")

    # --- Slide 10: Global Expansion (has OldCorp in title) ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "OldCorp Global Expansion Roadmap"
    tf10 = slide10.placeholders[1].text_frame
    tf10.paragraphs[0].text = "Phase 1 (2025): Southeast Asia launch - Vietnam, Thailand, Indonesia"
    add_paragraph(tf10, "Phase 2 (2026): Latin America - Brazil, Mexico, Colombia")
    add_paragraph(tf10, "Phase 3 (2027): Africa - Nigeria, Kenya, South Africa")

    # --- Slide 11: R&D Initiatives ---
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    slide11.shapes.title.text = "R&D Innovation Pipeline"
    tf11 = slide11.placeholders[1].text_frame
    tf11.paragraphs[0].text = "Autonomous delivery drone integration (Patent pending)"
    add_paragraph(tf11, "Blockchain-based supply chain verification")
    add_paragraph(tf11, "Digital twin simulation for warehouse design")
    add_paragraph(tf11, "Carbon footprint tracking and optimization")

    # --- Slide 12: Partnerships ---
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])
    slide12.shapes.title.text = "Strategic Partnerships"
    tf12 = slide12.placeholders[1].text_frame
    tf12.paragraphs[0].text = "Microsoft: Azure preferred logistics partner"
    add_paragraph(tf12, "Salesforce: CRM integration for sales logistics")
    add_paragraph(tf12, "SAP: ERP connectivity for enterprise clients")
    add_paragraph(tf12, "Maersk: Ocean freight data collaboration")

    # --- Slide 13: Sustainability Report ---
    slide13 = prs.slides.add_slide(prs.slide_layouts[1])
    slide13.shapes.title.text = "Sustainability & ESG Commitments"
    tf13 = slide13.placeholders[1].text_frame
    tf13.paragraphs[0].text = "50% carbon reduction target by 2030"
    add_paragraph(tf13, "100% renewable energy in all offices by 2027")
    add_paragraph(tf13, "EV fleet transition: 40% complete")
    add_paragraph(tf13, "Community investment: $12M in STEM education programs")

    # --- Slide 14: Competitive Analysis (has OldCorp in body) ---
    slide14 = prs.slides.add_slide(prs.slide_layouts[1])
    slide14.shapes.title.text = "Competitive Landscape"
    tf14 = slide14.placeholders[1].text_frame
    tf14.paragraphs[0].text = "OldCorp leads in AI-driven route optimization with 98.2% accuracy."
    add_paragraph(tf14, "Competitor A: Strong in warehouse automation but limited global reach.")
    add_paragraph(tf14, "Competitor B: Price-competitive but lacks real-time analytics.")
    add_paragraph(tf14, "Competitor C: Regional player in Europe with legacy technology stack.")

    # --- Slide 15: Risk Assessment ---
    slide15 = prs.slides.add_slide(prs.slide_layouts[1])
    slide15.shapes.title.text = "Risk Assessment & Mitigation"
    tf15 = slide15.placeholders[1].text_frame
    tf15.paragraphs[0].text = "Regulatory risk: GDPR compliance and data sovereignty requirements"
    add_paragraph(tf15, "Supply chain disruption: Diversified multi-region infrastructure")
    add_paragraph(tf15, "Talent acquisition: Competitive compensation and remote-first culture")
    add_paragraph(tf15, "Cybersecurity: SOC 2 Type II certified, annual pen testing")

    # --- Slide 16: Growth Projections ---
    slide16 = prs.slides.add_slide(prs.slide_layouts[1])
    slide16.shapes.title.text = "5-Year Growth Projections"
    tf16 = slide16.placeholders[1].text_frame
    tf16.paragraphs[0].text = "2026: $2.8B revenue (20% growth)"
    add_paragraph(tf16, "2027: $3.4B revenue (21% growth)")
    add_paragraph(tf16, "2028: $4.1B revenue (21% growth)")
    add_paragraph(tf16, "2029: $5.0B revenue (22% growth)")
    add_paragraph(tf16, "2030: $6.2B revenue (24% growth)")

    # --- Slide 17: Next Steps ---
    slide17 = prs.slides.add_slide(prs.slide_layouts[1])
    slide17.shapes.title.text = "Immediate Action Items"
    tf17 = slide17.placeholders[1].text_frame
    tf17.paragraphs[0].text = "Board approval for APAC expansion budget ($45M)"
    add_paragraph(tf17, "Finalize partnership agreement with Maersk by Q1 2026")
    add_paragraph(tf17, "Launch WarehouseIQ v3.0 beta program")
    add_paragraph(tf17, "Complete Series E fundraising ($200M target)")

    # --- Slide 18: Thank You ---
    slide18 = prs.slides.add_slide(prs.slide_layouts[0])
    slide18.shapes.title.text = "Thank You"
    slide18.placeholders[1].text = "Questions? Contact: strategy@oldcorp.com"

    # --- Add OldCorp to master slide footer ---
    # Access the slide master and add/modify footer placeholder
    slide_master = prs.slide_masters[0]
    # Add footer text to master via XML manipulation
    for shape in slide_master.shapes:
        if shape.has_text_frame:
            # Check for footer placeholder type
            ph = shape.placeholder_format
            if ph is not None and ph.type is not None:
                # placeholder type 12 = footer
                pass

    # Insert footer text into the slide master via XML
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # Enable footer on all slide layouts
    for slide_layout in prs.slide_layouts:
        for shape in slide_layout.placeholders:
            if shape.placeholder_format.idx == 11:  # footer placeholder
                shape.text = "OldCorp | Confidential"

    # Also set footer on the slide master
    for shape in slide_master.placeholders:
        if shape.placeholder_format.idx == 11:
            shape.text = "OldCorp | Confidential"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
