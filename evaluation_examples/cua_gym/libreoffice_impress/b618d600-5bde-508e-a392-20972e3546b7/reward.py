"""
Reward Script: Replace 'OldCorp' with 'NewCorp Technologies' across all slides
Task ID: impress_fix_092
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): No 'OldCorp' remains in any slide text
  Component 2 (0.3): 'NewCorp Technologies' present in expected slide locations
  Component 3 (0.3): Footer text in slide layouts replaced correctly
"""

import os
import zipfile

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_092'


def persist_app_state(domain):
    """Save any unsaved LibreOffice edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_all_slide_text(prs):
    """Extract all text from all slides, returning list of (slide_num, text) tuples."""
    results = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            texts = _extract_text_from_shape(shape)
            for t in texts:
                if t.strip():
                    results.append((i + 1, t))
    return results


def _extract_text_from_shape(shape):
    """Recursively extract text from a shape, including grouped shapes."""
    texts = []
    if hasattr(shape, 'text_frame') and hasattr(shape, 'text'):
        for para in shape.text_frame.paragraphs:
            full = ''.join(r.text or '' for r in para.runs)
            if full.strip():
                texts.append(full)
    if hasattr(shape, 'shapes'):
        for sub in shape.shapes:
            texts.extend(_extract_text_from_shape(sub))
    return texts


def get_layout_footer_texts(pptx_path):
    """Extract footer-like text from slide layouts via XML to catch OldCorp/NewCorp."""
    texts = []
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            for name in zf.namelist():
                if 'slideLayout' in name and name.endswith('.xml'):
                    content = zf.open(name).read().decode('utf-8')
                    texts.append((name, content))
    except Exception as e:
        print(f"ERROR: Cannot read ZIP: {e}")
    return texts


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: No 'OldCorp' remains in any slide text (0.4 points)
    # This checks that the old company name has been fully removed from slides.
    try:
        all_text = get_all_slide_text(prs)
        oldcorp_in_slides = [(snum, t) for snum, t in all_text if 'OldCorp' in t]
        if len(oldcorp_in_slides) == 0:
            print(f"PASS: Component 1 — No 'OldCorp' found in any slide text (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 — Found {len(oldcorp_in_slides)} occurrences of 'OldCorp' in slides:")
            for snum, t in oldcorp_in_slides:
                print(f"  Slide {snum}: {t[:80]}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: 'NewCorp Technologies' present in expected slide locations (0.3 points)
    # Task says OldCorp was in slides 1,2,3,5,7,10,14. After replacement, NewCorp Technologies
    # should appear in those slides.
    try:
        expected_slides = {1, 2, 3, 5, 7, 10, 14}
        all_text = get_all_slide_text(prs)
        slides_with_newcorp = set()
        for snum, t in all_text:
            if 'NewCorp Technologies' in t:
                slides_with_newcorp.add(snum)

        matched = expected_slides & slides_with_newcorp
        if len(matched) == len(expected_slides):
            print(f"PASS: Component 2 — 'NewCorp Technologies' found in all {len(expected_slides)} expected slides (0.3 pts)")
            total_score += 0.3
        else:
            missing = expected_slides - slides_with_newcorp
            print(f"FAIL: Component 2 — 'NewCorp Technologies' missing from slides: {sorted(missing)}")
            # Partial credit: proportional to how many slides have it
            partial = 0.3 * len(matched) / len(expected_slides)
            if partial > 0:
                print(f"  Partial credit: {partial:.2f} pts ({len(matched)}/{len(expected_slides)} slides)")
                total_score += partial
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Footer text in slide layouts replaced (0.3 points)
    # Initial has 'OldCorp | Confidential' in 11 slide layouts.
    # Golden should have 'NewCorp Technologies | Confidential' instead, with no 'OldCorp' remaining.
    try:
        layout_data = get_layout_footer_texts(file_path)
        layouts_with_oldcorp = 0
        layouts_with_newcorp = 0
        total_layouts_checked = 0

        for name, content in layout_data:
            total_layouts_checked += 1
            if 'OldCorp' in content:
                layouts_with_oldcorp += 1
            if 'NewCorp Technologies' in content:
                layouts_with_newcorp += 1

        no_old = (layouts_with_oldcorp == 0)
        has_new = (layouts_with_newcorp > 0)

        if no_old and has_new:
            print(f"PASS: Component 3 — No 'OldCorp' in layouts, 'NewCorp Technologies' found in {layouts_with_newcorp} layouts (0.3 pts)")
            total_score += 0.3
        else:
            if not no_old:
                print(f"FAIL: Component 3 — 'OldCorp' still present in {layouts_with_oldcorp} slide layouts")
            if not has_new:
                print(f"FAIL: Component 3 — 'NewCorp Technologies' not found in any slide layout")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist any unsaved changes before verification
persist_app_state("libreoffice_impress")

# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
