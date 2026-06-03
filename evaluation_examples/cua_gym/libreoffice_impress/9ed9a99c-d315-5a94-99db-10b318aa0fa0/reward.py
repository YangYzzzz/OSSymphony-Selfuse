"""
Reward Script: Scatter chart with trendline on slide 5
Task ID: impress_gf2_016
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.25): Scatter chart exists on slide 5
  - Component 2 (0.25): Chart title is 'Ad Spend vs. Sales Correlation'
  - Component 3 (0.30): 8 correct data points (X=Spend, Y=Sales)
  - Component 4 (0.20): Linear trendline present
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_016'

# Expected data points
EXPECTED_X = [10, 15, 20, 25, 30, 35, 40, 45]
EXPECTED_Y = [85, 110, 125, 160, 155, 195, 210, 240]


def persist_app_state(domain):
    """Save any unsaved LibreOffice state before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print(f"PERSIST: ctrl+s sent for {domain}")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide5 = prs.slides[4]  # 0-indexed

    # Find chart shape on slide 5
    chart_shape = None
    for shape in slide5.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Scatter chart exists on slide 5 (0.25 points)
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            # Verify it's a scatter/XY chart type
            chart_type_val = chart.chart_type
            # XY_SCATTER = -4169, XY_SCATTER_LINES = 74, etc.
            # Accept any XY scatter variant
            is_scatter = 'SCATTER' in str(chart_type_val).upper() or 'XY' in str(chart_type_val).upper()
            if is_scatter:
                print(f"PASS: Component 1 — Scatter chart found on slide 5, type={chart_type_val} (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 — Chart found but not scatter type, got: {chart_type_val}")
        else:
            print("FAIL: Component 1 — No chart found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    if chart_shape is None:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    chart = chart_shape.chart

    # Component 2: Chart title is 'Ad Spend vs. Sales Correlation' (0.25 points)
    try:
        if chart.has_title:
            title_text = chart.chart_title.text_frame.text.strip()
            if title_text == 'Ad Spend vs. Sales Correlation':
                print(f"PASS: Component 2 — Chart title matches exactly (0.25 pts)")
                total_score += 0.25
            else:
                # Partial: check if close
                if 'ad spend' in title_text.lower() and 'sales' in title_text.lower():
                    print(f"PARTIAL: Component 2 — Title contains key terms but not exact: '{title_text}' (0.10 pts)")
                    total_score += 0.10
                else:
                    print(f"FAIL: Component 2 — Expected 'Ad Spend vs. Sales Correlation', got: '{title_text}'")
        else:
            print("FAIL: Component 2 — Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: 8 correct data points (0.30 points)
    # Parse chart XML for X and Y values since python-pptx XySeries doesn't expose xValues easily
    try:
        c_ns = 'http://schemas.openxmlformats.org/drawingml/2006/chart'

        x_values = []
        y_values = []

        with zipfile.ZipFile(file_path, 'r') as zf:
            # Find chart XML files
            chart_files = [n for n in zf.namelist()
                           if n.startswith('ppt/charts/chart') and n.endswith('.xml')]
            for chart_file in chart_files:
                with zf.open(chart_file) as f:
                    root = ET.parse(f).getroot()

                # Check if this chart has scatterChart
                scatter_elems = list(root.iter(f'{{{c_ns}}}scatterChart'))
                if not scatter_elems:
                    continue

                # Extract xVal
                for xval_elem in root.iter(f'{{{c_ns}}}xVal'):
                    num_cache = xval_elem.find(f'.//{{{c_ns}}}numCache')
                    if num_cache is not None:
                        pts = num_cache.findall(f'{{{c_ns}}}pt')
                        x_values = [float(pt.find(f'{{{c_ns}}}v').text) for pt in pts]

                # Extract yVal
                for yval_elem in root.iter(f'{{{c_ns}}}yVal'):
                    num_cache = yval_elem.find(f'.//{{{c_ns}}}numCache')
                    if num_cache is not None:
                        pts = num_cache.findall(f'{{{c_ns}}}pt')
                        y_values = [float(pt.find(f'{{{c_ns}}}v').text) for pt in pts]

        if len(x_values) == 8 and len(y_values) == 8:
            # Check if values match expected (allow small tolerance for float)
            x_match = all(abs(a - b) < 0.5 for a, b in zip(x_values, EXPECTED_X))
            y_match = all(abs(a - b) < 0.5 for a, b in zip(y_values, EXPECTED_Y))

            if x_match and y_match:
                print(f"PASS: Component 3 — All 8 data points match (0.30 pts)")
                print(f"  X: {x_values}")
                print(f"  Y: {y_values}")
                total_score += 0.30
            elif x_match or y_match:
                print(f"PARTIAL: Component 3 — Only {'X' if x_match else 'Y'} values match (0.15 pts)")
                total_score += 0.15
            else:
                # Count individual matches
                x_correct = sum(1 for a, b in zip(x_values, EXPECTED_X) if abs(a - b) < 0.5)
                y_correct = sum(1 for a, b in zip(y_values, EXPECTED_Y) if abs(a - b) < 0.5)
                print(f"FAIL: Component 3 — Data mismatch. X correct: {x_correct}/8, Y correct: {y_correct}/8")
                print(f"  X actual: {x_values}")
                print(f"  Y actual: {y_values}")
        elif len(x_values) > 0 or len(y_values) > 0:
            print(f"PARTIAL: Component 3 — Wrong point count: X has {len(x_values)}, Y has {len(y_values)} (expected 8 each)")
            # Give small credit for having some data
            if len(x_values) > 0 and len(y_values) > 0:
                total_score += 0.10
        else:
            print("FAIL: Component 3 — No data points found in chart")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Linear trendline present (0.20 points)
    try:
        trendline_found = False
        is_linear = False

        with zipfile.ZipFile(file_path, 'r') as zf:
            chart_files = [n for n in zf.namelist()
                           if n.startswith('ppt/charts/chart') and n.endswith('.xml')]
            for chart_file in chart_files:
                with zf.open(chart_file) as f:
                    root = ET.parse(f).getroot()

                for trendline in root.iter(f'{{{c_ns}}}trendline'):
                    trendline_found = True
                    ttype = trendline.find(f'{{{c_ns}}}trendlineType')
                    if ttype is not None and ttype.get('val') == 'linear':
                        is_linear = True

        if is_linear:
            print("PASS: Component 4 — Linear trendline found (0.20 pts)")
            total_score += 0.20
        elif trendline_found:
            print("PARTIAL: Component 4 — Trendline found but not linear type (0.10 pts)")
            total_score += 0.10
        else:
            print("FAIL: Component 4 — No trendline found in chart")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
