"""
Initial Setup: Marketing Analytics presentation with 7 slides, slide 5 has title
'Correlation Analysis' and empty content area. No charts exist.
Task ID: impress_gf2_016
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_016'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=None, alignment=None):
    """Helper to add a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=12):
    """Add a bullet list text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(font_size)
    return txBox


def create_initial():
    prs = Presentation()
    # Standard 16:9 slide dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Marketing Analytics Report"
    slide1.placeholders[1].text = "Q1 2025 Performance Review\nPrepared by Analytics Team"

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_text_box(slide2, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Executive Summary", font_size=28, bold=True,
                 color=RGBColor(0x1F, 0x3A, 0x6E))
    add_bullet_list(slide2, Inches(0.8), Inches(1.6), Inches(11), Inches(5), [
        "Total marketing spend increased 18% year-over-year to $2.4M across all channels",
        "Digital advertising contributed 62% of total qualified leads, up from 48% in Q4 2024",
        "Customer acquisition cost decreased from $127 to $98 per converted lead",
        "Social media engagement rates improved 34% following the new content strategy rollout",
        "Email campaign open rates averaged 24.3%, exceeding the industry benchmark of 21.5%",
        "Brand awareness surveys indicate a 12-point increase in unaided recall among target demographics",
    ], font_size=14)

    # --- Slide 3: Campaign Overview (Table) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide3, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Campaign Overview", font_size=28, bold=True,
                 color=RGBColor(0x1F, 0x3A, 0x6E))

    rows, cols = 7, 5
    tbl_shape = slide3.shapes.add_table(rows, cols, Inches(0.8), Inches(1.6),
                                         Inches(11), Inches(4.5))
    table = tbl_shape.table
    headers = ["Campaign", "Channel", "Budget", "Leads", "Conv. Rate"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1F, 0x3A, 0x6E)

    campaigns = [
        ["Spring Launch", "Google Ads", "$185,000", "2,340", "3.8%"],
        ["Brand Awareness", "Facebook", "$142,500", "1,870", "2.1%"],
        ["Product Demo", "LinkedIn", "$98,000", "1,120", "5.2%"],
        ["Retargeting", "Display Network", "$67,500", "890", "7.4%"],
        ["Email Nurture", "Email", "$45,000", "1,560", "4.6%"],
        ["Influencer Collab", "Instagram", "$112,000", "2,100", "3.1%"],
    ]
    for r, row_data in enumerate(campaigns, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 4: Channel Performance ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide4, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Channel Performance", font_size=28, bold=True,
                 color=RGBColor(0x1F, 0x3A, 0x6E))
    add_bullet_list(slide4, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5), [
        "Google Ads: 42% of total spend, highest lead volume",
        "Facebook: Strong reach but lower conversion rates",
        "LinkedIn: Best quality leads with 5.2% conversion",
        "Instagram: Growing channel, 23% engagement increase",
        "Email: Lowest cost per acquisition at $29/lead",
        "Display Network: Highest conversion rate at 7.4%",
    ], font_size=13)
    add_bullet_list(slide4, Inches(6.8), Inches(1.6), Inches(5.5), Inches(5), [
        "Top Performer: Display retargeting ($76 CPA)",
        "Most Improved: Instagram (+23% YoY engagement)",
        "Highest ROI: Email marketing (312% return)",
        "Largest Budget: Google Ads ($185K allocated)",
        "Best Awareness: Facebook (4.2M impressions)",
    ], font_size=13)

    # --- Slide 5: Correlation Analysis (empty content area, no charts) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide5, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Correlation Analysis", font_size=28, bold=True,
                 color=RGBColor(0x1F, 0x3A, 0x6E))
    # Empty white content area placeholder (no chart!)
    content_box = slide5.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0.8), Inches(1.6), Inches(11), Inches(5.2)
    )
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    content_box.line.color.rgb = RGBColor(0xD9, 0xD9, 0xD9)

    # --- Slide 6: Budget Recommendations ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide6, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Budget Recommendations", font_size=28, bold=True,
                 color=RGBColor(0x1F, 0x3A, 0x6E))
    add_bullet_list(slide6, Inches(0.8), Inches(1.6), Inches(11), Inches(5), [
        "Increase display retargeting budget by 25% based on strong conversion metrics",
        "Reallocate $30K from Facebook brand awareness to LinkedIn lead generation",
        "Expand email automation sequences to capture mid-funnel prospects more effectively",
        "Test new TikTok advertising channel with $50K pilot budget in Q2",
        "Maintain Google Ads spend but optimize keyword targeting for lower CPA",
        "Invest in marketing attribution tooling to improve cross-channel measurement",
    ], font_size=14)

    # --- Slide 7: Next Steps ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide7, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Next Steps", font_size=28, bold=True,
                 color=RGBColor(0x1F, 0x3A, 0x6E))
    add_bullet_list(slide7, Inches(0.8), Inches(1.6), Inches(11), Inches(5), [
        "Schedule Q2 campaign planning workshop with creative and media teams",
        "Finalize vendor selection for new marketing attribution platform by April 15",
        "Launch A/B testing framework for landing page optimization across all channels",
        "Complete competitive benchmarking analysis and present findings to leadership",
        "Develop integrated reporting dashboard connecting spend to revenue metrics",
    ], font_size=14)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
