"""
Initial Setup: Create a 15-slide Brand Guidelines presentation and open in Impress
Task ID: impress_rp_035
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_035'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
EXPORTS_DIR = f'{WORKDIR}/exports'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=None, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    """Helper to add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullet_slide(slide, title_text, bullets, title_color=RGBColor(0x1B, 0x3A, 0x5C)):
    """Add a title and bullet list to a slide."""
    # Title
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(1),
                 title_text, font_size=28, bold=True, color=title_color)
    # Bullets
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BRAND_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
    BRAND_TEAL = RGBColor(0x00, 0x96, 0x88)
    BRAND_GRAY = RGBColor(0x60, 0x60, 0x60)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    # ── Slide 1: Title Slide ──
    s1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill = s1.background.fill
    fill.solid()
    fill.fore_color.rgb = BRAND_BLUE
    add_text_box(s1, Inches(1.5), Inches(1.5), Inches(10), Inches(1.5),
                 "Meridian Technologies", font_size=44, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(s1, Inches(1.5), Inches(3.2), Inches(10), Inches(1),
                 "Brand Guidelines 2025", font_size=28, bold=False, color=BRAND_TEAL,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(s1, Inches(1.5), Inches(4.5), Inches(10), Inches(0.8),
                 "Confidential — For Internal Use Only", font_size=14, bold=False,
                 color=WHITE, alignment=PP_ALIGN.CENTER)

    # ── Slide 2: Table of Contents ──
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(s2, Inches(0.8), Inches(0.4), Inches(8), Inches(1),
                 "Table of Contents", font_size=28, bold=True, color=BRAND_BLUE)
    toc_items = [
        "1. Brand Overview & Mission",
        "2. Logo Usage & Clear Space",
        "3. Primary Color Palette",
        "4. Secondary Color Palette",
        "5. Typography Standards",
        "6. Voice & Tone",
        "7. Photography Guidelines",
        "8. Iconography Standards",
        "9. Digital Media Specifications",
        "10. Print Media Specifications",
        "11. Email Signature & Templates",
        "12. Social Media Guidelines",
        "13. Do's and Don'ts"
    ]
    txBox = s2.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(toc_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.space_after = Pt(6)
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(15)
        run.font.color.rgb = BRAND_GRAY

    # ── Slide 3: Brand Overview ──
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(s3, "Brand Overview & Mission", [
        "Meridian Technologies was founded in 2012 with the vision of making enterprise software intuitive.",
        "Our mission: Empower organizations to achieve operational excellence through intelligent automation.",
        "Core values: Innovation, Integrity, Inclusivity, Impact.",
        "We serve over 4,500 enterprise clients across 38 countries.",
        "Annual revenue exceeded $2.3 billion in fiscal year 2024.",
    ])

    # ── Slide 4: Logo Usage ──
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(s4, "Logo Usage & Clear Space", [
        "The primary logo must always appear on white or light gray (#F5F5F5) backgrounds.",
        "Minimum clear space: 1.5x the height of the logo mark on all sides.",
        "Never stretch, rotate, or apply drop shadows to the logo.",
        "Monochrome version is available for single-color printing.",
        "Minimum reproduction size: 24mm width for print, 120px for digital.",
        "Logo lockup must include the wordmark below 200px width.",
    ])

    # ── Slide 5: Primary Colors ──
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(s5, Inches(0.8), Inches(0.4), Inches(8), Inches(1),
                 "Primary Color Palette", font_size=28, bold=True, color=BRAND_BLUE)
    colors_data = [
        ("Meridian Blue", "#1B3A5C", "RGB(27, 58, 92)", "Primary brand color — headers, CTAs"),
        ("Meridian Teal", "#009688", "RGB(0, 150, 136)", "Accent — links, highlights, icons"),
        ("Charcoal", "#333333", "RGB(51, 51, 51)", "Body text — paragraphs, captions"),
        ("Cloud White", "#FAFAFA", "RGB(250, 250, 250)", "Backgrounds — cards, containers"),
    ]
    table_shape = s5.shapes.add_table(len(colors_data) + 1, 4,
                                       Inches(0.8), Inches(1.8),
                                       Inches(11), Inches(3.5))
    table = table_shape.table
    headers = ["Name", "Hex", "RGB", "Usage"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = WHITE
        from pptx.oxml.ns import qn
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '1B3A5C'})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)
    for r, (name, hexval, rgb, usage) in enumerate(colors_data, 1):
        for c, val in enumerate([name, hexval, rgb, usage]):
            cell = table.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(13)
                run.font.color.rgb = BRAND_GRAY

    # ── Slide 6: Secondary Colors ──
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(s6, "Secondary Color Palette", [
        "Sunset Orange (#E8703A) — Warnings, urgent alerts, promotional badges.",
        "Forest Green (#2E7D32) — Success states, confirmation messages, eco-branding.",
        "Lavender (#7E57C2) — Innovation sections, premium features, creative assets.",
        "Soft Gold (#FFB74D) — Awards, highlights, achievement badges.",
        "Slate Gray (#78909C) — Borders, dividers, secondary text.",
    ])

    # ── Slide 7: Typography Standards ──
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(s7, "Typography Standards", [
        "Primary Typeface: Inter (Headings — Bold, 28–44pt; Subheadings — SemiBold, 20–24pt).",
        "Secondary Typeface: Source Sans Pro (Body text — Regular, 14–16pt).",
        "Monospace: Fira Code (Code snippets, data tables — Regular, 12–14pt).",
        "Line height: 1.5x for body text, 1.2x for headings.",
        "Maximum two typeface families per page.",
        "Never use decorative or script fonts in any official material.",
    ])

    # ── Slide 8: Voice & Tone ──
    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(s8, "Voice & Tone", [
        "Professional but approachable — avoid jargon unless addressing a technical audience.",
        "Active voice preferred: 'We build solutions' over 'Solutions are built by us.'",
        "Contractions are acceptable in marketing and blog content (e.g., 'you'll', 'we're').",
        "Avoid superlatives without evidence: 'leading' requires a cited source.",
        "Customer success stories should lead with the outcome, not the product.",
    ])

    # ── Slide 9: Photography Guidelines ──
    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(s9, "Photography Guidelines", [
        "Use authentic, high-resolution images (minimum 300 DPI for print, 72 DPI for web).",
        "Preferred subjects: diverse teams collaborating, modern workspaces, technology in use.",
        "Avoid stock photos with artificial poses or overused compositions.",
        "Color grading should maintain warm neutrals — avoid heavy blue or orange filters.",
        "All images must pass accessibility contrast checks (WCAG 2.1 AA).",
    ])

    # ── Slide 10: Iconography Standards ──
    s10 = prs.slides.add_slide(prs.slide_layouts[6])
    fill = s10.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    add_text_box(s10, Inches(0.8), Inches(0.4), Inches(8), Inches(1),
                 "Iconography Standards", font_size=28, bold=True, color=BRAND_BLUE)
    icon_details = [
        "Style: Outlined icons with 2px stroke weight, rounded corners (2px radius).",
        "Grid: All icons designed on a 24x24px grid with 2px padding.",
        "Colors: Single-color icons only — use Meridian Blue or Charcoal.",
        "Sizing: 16px (inline), 24px (navigation), 32px (feature highlights), 48px (hero sections).",
        "File formats: SVG for web, PDF for print, PNG @2x for email templates.",
        "Consistency: Use the Meridian Icon Library (v3.2). Custom icons require design team approval.",
    ]
    txBox = s10.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(icon_details):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.space_after = Pt(10)
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(16)
        run.font.color.rgb = BRAND_GRAY

    # ── Slide 11: Digital Media Specifications ──
    s11 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(s11, "Digital Media Specifications", [
        "Website banner: 1920x600px, JPG (quality 85%) or WebP.",
        "Social media profile picture: 400x400px, PNG with transparent background.",
        "LinkedIn cover: 1584x396px; Twitter/X header: 1500x500px.",
        "Email header: 600px wide, maximum 200KB file size.",
        "App store screenshots: 1284x2778px (iPhone), 2048x2732px (iPad).",
    ])

    # ── Slide 12: Print Media Specifications ──
    s12 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(s12, "Print Media Specifications", [
        "Business cards: 3.5 x 2 inches, 300 DPI, CMYK color mode.",
        "Letterhead: A4 (210x297mm), 15mm margins, brand header at top.",
        "Brochures: Tri-fold (8.5 x 11in) or bi-fold (11 x 17in), 3mm bleed.",
        "Trade show banners: 33 x 81 inches, 150 DPI minimum.",
        "All print materials must include Pantone color codes (PMS 302 C for Meridian Blue).",
    ])

    # ── Slide 13: Email Signature ──
    s13 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(s13, "Email Signature & Templates", [
        "Standard format: Full Name | Title | Department.",
        "Phone and email on the second line, separated by a pipe character.",
        "Logo placement: left-aligned, 120px wide, linked to company website.",
        "Social icons: LinkedIn and Twitter/X only, 16px, Meridian Blue.",
        "Email templates use the Inter font stack with Source Sans Pro fallback.",
        "Maximum signature height: 120px including logo.",
    ])

    # ── Slide 14: Social Media ──
    s14 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(s14, "Social Media Guidelines", [
        "Hashtag policy: Use #MeridianTech and campaign-specific tags (max 5 per post).",
        "Response SLA: Customer inquiries within 2 hours during business hours.",
        "Repost/share only from verified partner accounts.",
        "Video content: 16:9 aspect ratio, captions required, max 60 seconds for stories.",
        "Approval workflow: All posts require Marketing Manager sign-off before publishing.",
    ])

    # ── Slide 15: Do's and Don'ts ──
    s15 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(s15, Inches(0.8), Inches(0.4), Inches(8), Inches(1),
                 "Do's and Don'ts", font_size=28, bold=True, color=BRAND_BLUE)
    # Do's column
    add_text_box(s15, Inches(0.8), Inches(1.5), Inches(5), Inches(0.6),
                 "DO", font_size=20, bold=True, color=RGBColor(0x2E, 0x7D, 0x32))
    dos = [
        "Use approved templates for all external communications.",
        "Maintain minimum clear space around the logo.",
        "Reference the color palette for all digital and print assets.",
        "Get written approval before creating co-branded materials.",
    ]
    txBox = s15.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(5), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(dos):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"✓ {item}"
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0x2E, 0x7D, 0x32)

    # Don'ts column
    add_text_box(s15, Inches(7), Inches(1.5), Inches(5), Inches(0.6),
                 "DON'T", font_size=20, bold=True, color=RGBColor(0xC6, 0x28, 0x28))
    donts = [
        "Alter the logo proportions or apply unapproved effects.",
        "Use brand colors outside the defined palette.",
        "Publish content without going through the approval process.",
        "Use competitor comparisons without Legal review.",
    ]
    txBox2 = s15.shapes.add_textbox(Inches(7), Inches(2.2), Inches(5), Inches(4))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(donts):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"✗ {item}"
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0xC6, 0x28, 0x28)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Presentation has {len(prs.slides)} slides')

    # Create empty exports directory
    os.makedirs(EXPORTS_DIR, exist_ok=True)
    print(f'Exports directory created: {EXPORTS_DIR}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_presentation()
