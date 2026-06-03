"""
Reward Script: Insert horizontal line divider on slide 2
Task ID: impress_gf3_026
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25) — Line shape exists on slide 2
  Component 2 (0.25) — Line spans full slide width (approx 0 to 25.4 cm)
  Component 3 (0.20) — Line positioned at approximately y=4 cm from top
  Component 4 (0.15) — Line is 2pt thick
  Component 5 (0.15) — Line color is #333333 and style is solid
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_026'


def persist_app_state(domain):
    """Try to save any unsaved GUI edits via Ctrl+S."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def find_line_shapes_on_slide(slide):
    """Find all line/connector shapes on a slide.
    A line shape has shape_type LINE (9) or is a cxnSp element with prstGeom prst='line'.
    """
    lines = []
    for shape in slide.shapes:
        is_line = False
        # Check shape_type
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            is_line = True
        # Also check XML for connector shapes with line geometry
        el = shape._element
        tag = el.tag
        if 'cxnSp' in tag:
            geom = el.find('.//' + qn('a:prstGeom'))
            if geom is not None and geom.get('prst') == 'line':
                is_line = True
        if is_line:
            lines.append(shape)
    return lines


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file has at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 2")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]  # 0-indexed, slide 2
    slide_width_emu = prs.slide_width  # expected ~9144000

    # Find line shapes on slide 2
    line_shapes = find_line_shapes_on_slide(slide2)

    # Component 1: Line shape exists on slide 2 (0.25 points)
    try:
        if len(line_shapes) > 0:
            print(f"PASS: Component 1 — Found {len(line_shapes)} line shape(s) on slide 2 (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — No line shapes found on slide 2")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    if len(line_shapes) == 0:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Use the first (or best matching) line shape for remaining checks
    # If multiple lines, find the one closest to y=4cm
    target_y_emu = 1440000  # 4 cm in EMU
    best_line = min(line_shapes, key=lambda s: abs(s.top - target_y_emu))

    # Component 2: Line spans full slide width (0.25 points)
    # The line should go from approximately x=0 to the full slide width
    try:
        line_left = best_line.left
        line_width = best_line.width
        line_right = line_left + line_width
        # Tolerance: within 10% of slide width
        width_ratio = line_width / slide_width_emu
        starts_near_left = line_left <= slide_width_emu * 0.10  # starts within 10% of left edge
        width_is_full = width_ratio >= 0.85  # covers at least 85% of slide width

        if starts_near_left and width_is_full:
            print(f"PASS: Component 2 — Line spans {width_ratio*100:.1f}% of slide width, "
                  f"left={line_left/360000:.2f}cm (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 — Line width={line_width/360000:.2f}cm "
                  f"({width_ratio*100:.1f}% of slide), left={line_left/360000:.2f}cm")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Line positioned at approximately y=4 cm from top (0.20 points)
    try:
        line_top = best_line.top
        line_top_cm = line_top / 360000
        # Tolerance: within 1.0 cm of target (4 cm)
        if abs(line_top_cm - 4.0) <= 1.0:
            print(f"PASS: Component 3 — Line y-position = {line_top_cm:.2f} cm (target ~4.0 cm) (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 3 — Line y-position = {line_top_cm:.2f} cm, expected ~4.0 cm")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Line is 2pt thick (0.15 points)
    try:
        el = best_line._element
        ln = el.find('.//' + qn('a:ln'))
        if ln is not None:
            w_attr = ln.get('w')
            if w_attr is not None:
                width_pt = int(w_attr) / 12700.0
                # Tolerance: within 0.5pt of 2pt
                if abs(width_pt - 2.0) <= 0.5:
                    print(f"PASS: Component 4 — Line width = {width_pt:.2f} pt (target 2pt) (0.15 pts)")
                    total_score += 0.15
                else:
                    print(f"FAIL: Component 4 — Line width = {width_pt:.2f} pt, expected 2pt")
            else:
                print(f"FAIL: Component 4 — No width attribute on line element")
        else:
            print(f"FAIL: Component 4 — No line properties (a:ln) found")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Line color is #333333 and style is solid (0.15 points)
    try:
        el = best_line._element
        ln = el.find('.//' + qn('a:ln'))
        color_ok = False
        style_ok = False

        if ln is not None:
            # Check color
            solid_fill = ln.find(qn('a:solidFill'))
            if solid_fill is not None:
                srgb = solid_fill.find(qn('a:srgbClr'))
                if srgb is not None:
                    color_val = srgb.get('val', '').upper()
                    if color_val == '333333':
                        color_ok = True
                        print(f"  Color check: #{color_val} — matches #333333")
                    else:
                        print(f"  Color check: #{color_val} — expected #333333")
                else:
                    print(f"  Color check: no srgbClr found in solidFill")
            else:
                print(f"  Color check: no solidFill in line properties")

            # Check dash style — solid means either prstDash val="solid" or no dash element
            prst_dash = ln.find(qn('a:prstDash'))
            if prst_dash is not None:
                dash_val = prst_dash.get('val', '')
                if dash_val == 'solid':
                    style_ok = True
                    print(f"  Style check: {dash_val} — solid")
                else:
                    print(f"  Style check: {dash_val} — expected solid")
            else:
                # No dash element means solid by default
                style_ok = True
                print(f"  Style check: no dash element — defaults to solid")

        if color_ok and style_ok:
            print(f"PASS: Component 5 — Color #333333, solid style (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 5 — color_ok={color_ok}, style_ok={style_ok}")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entrypoint
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
