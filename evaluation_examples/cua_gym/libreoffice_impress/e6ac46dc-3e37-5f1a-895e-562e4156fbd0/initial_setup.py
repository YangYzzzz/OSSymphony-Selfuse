"""
Initial Setup: Insert a horizontal line shape on slide 2 of a presentation.
Task ID: impress_gf3_026
Domain: libreoffice_impress

Creates a 10-slide presentation with realistic content. Slide 2 has a title
and content area with a visible gap between them. No decorative lines or
shapes exist on slide 2.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_026'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=None):
    """Helper to add a styled textbox to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    # Standard 16:9 slide dimensions
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(19.05)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Quarterly Performance Review"
    slide1.placeholders[1].text = "FY2025 Q1 — Global Operations Division"

    # ---- Slide 2: Overview (target slide - title + content with gap) ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    # Title area at top
    add_textbox(slide2, Cm(1.5), Cm(0.8), Cm(22), Cm(2.5),
                "Executive Summary", font_size=28, bold=True,
                color=RGBColor(0x1A, 0x3C, 0x6E))
    # Content area below (with a gap from the title, starting at ~5.5 cm)
    content_lines = [
        "Revenue grew 12.3% year-over-year, reaching $48.7M in Q1 2025.",
        "Customer acquisition cost decreased by 8% compared to Q4 2024.",
        "Employee satisfaction scores improved across all departments.",
        "Three new product lines launched ahead of schedule.",
        "International expansion into Southeast Asian markets on track.",
    ]
    y_start = Cm(5.5)
    for i, line in enumerate(content_lines):
        add_textbox(slide2, Cm(2.0), y_start + Cm(i * 2.2), Cm(21), Cm(1.8),
                    line, font_size=14, color=RGBColor(0x33, 0x33, 0x33))

    # ---- Slide 3: Revenue Breakdown ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide3, Cm(1.5), Cm(0.8), Cm(22), Cm(2.5),
                "Revenue Breakdown by Region", font_size=26, bold=True,
                color=RGBColor(0x1A, 0x3C, 0x6E))
    regions = [
        ("North America", "$22.4M", "+15.1%"),
        ("Europe", "$14.2M", "+9.8%"),
        ("Asia-Pacific", "$8.3M", "+18.6%"),
        ("Latin America", "$3.8M", "+5.2%"),
    ]
    table_shape = slide3.shapes.add_table(5, 3, Cm(3), Cm(4), Cm(19), Cm(8))
    table = table_shape.table
    headers = ["Region", "Revenue", "YoY Growth"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, (region, rev, growth) in enumerate(regions, 1):
        table.cell(r, 0).text = region
        table.cell(r, 1).text = rev
        table.cell(r, 2).text = growth

    # ---- Slide 4: Customer Metrics ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide4, Cm(1.5), Cm(0.8), Cm(22), Cm(2.5),
                "Customer Metrics", font_size=26, bold=True,
                color=RGBColor(0x1A, 0x3C, 0x6E))
    metrics = [
        "Total Active Customers: 14,832 (+1,247 net new)",
        "Monthly Active Users: 9.2M (62% engagement rate)",
        "Net Promoter Score: 72 (up from 68 in Q4)",
        "Average Revenue Per User: $328/month",
        "Churn Rate: 2.1% (lowest in 3 years)",
    ]
    for i, m in enumerate(metrics):
        add_textbox(slide4, Cm(2.0), Cm(4.0) + Cm(i * 2.5), Cm(21), Cm(2.0),
                    m, font_size=16, color=RGBColor(0x44, 0x44, 0x44))

    # ---- Slide 5: Product Launch Timeline ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide5, Cm(1.5), Cm(0.8), Cm(22), Cm(2.5),
                "Product Launch Timeline", font_size=26, bold=True,
                color=RGBColor(0x1A, 0x3C, 0x6E))
    launches = [
        "Jan 15 — CloudSync Pro 3.0 released (Enterprise tier)",
        "Feb 8 — Mobile Analytics Dashboard v2.1 shipped",
        "Mar 1 — DataVault Security Suite entered beta",
        "Mar 22 — API Gateway 4.0 public preview launched",
    ]
    for i, l in enumerate(launches):
        add_textbox(slide5, Cm(2.5), Cm(4.5) + Cm(i * 2.8), Cm(20), Cm(2.0),
                    l, font_size=15, color=RGBColor(0x33, 0x33, 0x33))

    # ---- Slide 6: Team Performance ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide6, Cm(1.5), Cm(0.8), Cm(22), Cm(2.5),
                "Team Performance Highlights", font_size=26, bold=True,
                color=RGBColor(0x1A, 0x3C, 0x6E))
    teams = [
        ("Engineering", "Delivered 94% of sprint commitments"),
        ("Sales", "Exceeded quota by 17% ($8.2M over target)"),
        ("Marketing", "Campaign ROI improved 23% with new attribution model"),
        ("Support", "First response time reduced to 12 minutes average"),
    ]
    for i, (team, desc) in enumerate(teams):
        add_textbox(slide6, Cm(2.0), Cm(4.0) + Cm(i * 3.0), Cm(6), Cm(2.0),
                    team, font_size=18, bold=True, color=RGBColor(0x1A, 0x3C, 0x6E))
        add_textbox(slide6, Cm(8.5), Cm(4.0) + Cm(i * 3.0), Cm(15), Cm(2.0),
                    desc, font_size=14, color=RGBColor(0x55, 0x55, 0x55))

    # ---- Slide 7: Budget Overview ----
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide7, Cm(1.5), Cm(0.8), Cm(22), Cm(2.5),
                "Budget Overview", font_size=26, bold=True,
                color=RGBColor(0x1A, 0x3C, 0x6E))
    budget_items = [
        ("Personnel", "$18.5M", "38%"),
        ("Infrastructure", "$9.2M", "19%"),
        ("Marketing", "$7.1M", "15%"),
        ("R&D", "$11.3M", "23%"),
        ("Operations", "$2.6M", "5%"),
    ]
    tbl = slide7.shapes.add_table(6, 3, Cm(3), Cm(4), Cm(19), Cm(10))
    t = tbl.table
    for c, h in enumerate(["Category", "Amount", "% of Total"]):
        t.cell(0, c).text = h
        for run in t.cell(0, c).text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, (cat, amt, pct) in enumerate(budget_items, 1):
        t.cell(r, 0).text = cat
        t.cell(r, 1).text = amt
        t.cell(r, 2).text = pct

    # ---- Slide 8: Risk Assessment ----
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide8, Cm(1.5), Cm(0.8), Cm(22), Cm(2.5),
                "Risk Assessment", font_size=26, bold=True,
                color=RGBColor(0x1A, 0x3C, 0x6E))
    risks = [
        "Supply chain disruptions in APAC — Mitigation: dual-source strategy",
        "Talent retention in engineering — Mitigation: revised comp packages",
        "Regulatory changes in EU data privacy — Mitigation: compliance roadmap",
        "Currency fluctuation impact on international revenue",
    ]
    for i, risk in enumerate(risks):
        add_textbox(slide8, Cm(2.0), Cm(4.0) + Cm(i * 2.8), Cm(21), Cm(2.2),
                    risk, font_size=14, color=RGBColor(0x44, 0x44, 0x44))

    # ---- Slide 9: Q2 Priorities ----
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide9, Cm(1.5), Cm(0.8), Cm(22), Cm(2.5),
                "Q2 2025 Strategic Priorities", font_size=26, bold=True,
                color=RGBColor(0x1A, 0x3C, 0x6E))
    priorities = [
        "1. Complete Southeast Asia market entry (Malaysia, Vietnam)",
        "2. Launch DataVault Security Suite to general availability",
        "3. Achieve 10M monthly active users milestone",
        "4. Reduce infrastructure costs by 12% through optimization",
        "5. Hire 45 additional engineers for platform team",
    ]
    for i, p in enumerate(priorities):
        add_textbox(slide9, Cm(2.0), Cm(4.0) + Cm(i * 2.5), Cm(21), Cm(2.0),
                    p, font_size=15, color=RGBColor(0x33, 0x33, 0x33))

    # ---- Slide 10: Thank You ----
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide10, Cm(4), Cm(6), Cm(17), Cm(4),
                "Thank You", font_size=40, bold=True,
                color=RGBColor(0x1A, 0x3C, 0x6E),
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide10, Cm(4), Cm(10), Cm(17), Cm(3),
                "Questions & Discussion", font_size=22,
                color=RGBColor(0x66, 0x66, 0x66),
                alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
