"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 46 looks messy—every shape, picture, and text box shifts a few pixels left or right. In LibreOffice Impress, what’s the fastest way to grab all the objects on that single slide and align them to the exact horizontal center of the slide (Align ➜ Center)?
Generated: 2025-09-10 14:13:12
Status: success
Model: azure-o3
Total Steps: 6
"""

import os
from pptx import Presentation


def verify_horizontal_center_alignment(file_path: str, slide_number: int = 46, tolerance_px: int = 5) -> float:
    """Verify that every measurable shape on `slide_number` is horizontally
    centered (Align ➜ Center) within `tolerance_px` pixels.

    Returns a progressive score between 0.0 and 1.0 based on the proportion
    of shapes correctly aligned. A score of 1.0 means every applicable shape
    is within the tolerance of the exact horizontal center.
    """

    print(f"Loading presentation: {file_path}")

    # --- Basic file sanity check (no points awarded) ---
    if not os.path.exists(file_path):
        print("✗ File not found")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Unable to load presentation: {e}")
        return 0.0  # loading failure = task not accomplished

    # Convert pixel tolerance to EMU (English Metric Units)
    PIXEL_TO_EMU = 9525  # 1 px ≈ 9525 EMU at 96 DPI
    tolerance_emu = tolerance_px * PIXEL_TO_EMU

    slide_index = slide_number - 1  # zero-based index
    if slide_index >= len(prs.slides):
        print(f"✗ Slide {slide_number} does not exist (presentation has {len(prs.slides)} slides)")
        return 0.0

    slide = prs.slides[slide_index]
    slide_center = prs.slide_width / 2  # horizontal midpoint of the slide

    align_results = []  # True/False per measurable shape

    print(f"Verifying horizontal alignment on slide {slide_number} (center = {slide_center} EMU, tolerance = ±{tolerance_emu} EMU)…")

    for shape in slide.shapes:
        # Only evaluate shapes that expose geometric properties and have non-zero width
        if not hasattr(shape, "left") or not hasattr(shape, "width"):
            continue
        if shape.width == 0:
            continue

        cx = shape.left + shape.width / 2  # shape horizontal center
        delta = abs(cx - slide_center)
        aligned = delta <= tolerance_emu
        align_results.append(aligned)

        print(
            f"  Shape type {shape.shape_type:>2}: center_delta = {delta} EMU "
            f"({'✓' if aligned else '✗'})"
        )

    # If no measurable shapes exist, task cannot be evaluated → 0.0
    if not align_results:
        print("✗ No measurable shapes found on the slide; cannot verify alignment")
        return 0.0

    aligned_count = sum(align_results)
    total_count = len(align_results)
    score = aligned_count / total_count  # progressive scoring

    print(f"Aligned shapes: {aligned_count}/{total_count}")
    print(f"Calculated score: {score}")

    return score


# ----------------- MAIN EXECUTION -----------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/slide_46_looks_messyevery_shape_picture_and_text_box_shifts_a_few_pixels_left_or_right_in_libreoffic_golden.pptx"
    reward = verify_horizontal_center_alignment(FILE_PATH)
    print(f"REWARD: {reward}")
