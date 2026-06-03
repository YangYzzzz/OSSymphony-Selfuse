"""
Initial Setup: Lock the background rectangle on slide 1
Task ID: impress_objects_049
Domain: libreoffice_impress

Creates a presentation file with:
- Slide 1: Large background rectangle (unlocked) + title + content shapes
- Slide 2: Agenda slide with bullet points
- Slide 3: Content slide with text and a table

The background rectangle on slide 1 must NOT be locked (task is to lock it).
"""

import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import lxml.etree as etree

WORKDIR = '/home/user'
TASK_ID = 'impress_objects_049'
OUTPUT = f'{WORKDIR}/{TASK_ID}_initial.pptx'
DESKTOP_OUTPUT = f'{WORKDIR}/Desktop/template_base.pptx'


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ─── Slide 1: Title slide with background rectangle ───
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    # Background rectangle — large, covers most of the slide
    # This is the shape that the task asks to lock
    bg_rect = slide1.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0),
        Inches(13.33), Inches(7.5)
    )
    bg_rect.name = 'BackgroundRect'
    # Style: dark navy fill
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x4A)
    bg_rect.line.color.rgb = RGBColor(0x1A, 0x2A, 0x4A)
    bg_rect.line.width = Pt(0)

    # Decorative accent bar
    accent_bar = slide1.shapes.add_shape(
        1,
        Inches(0), Inches(6.5),
        Inches(13.33), Inches(0.15)
    )
    accent_bar.name = 'AccentBar'
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = RGBColor(0xFF, 0xA5, 0x00)
    accent_bar.line.fill.background()

    # Title text box
    title_box = slide1.shapes.add_textbox(
        Inches(1.5), Inches(2.0),
        Inches(10.0), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = 'Q4 Business Review'
    run.font.name = 'Calibri'
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Subtitle text box
    sub_box = slide1.shapes.add_textbox(
        Inches(2.0), Inches(3.8),
        Inches(9.0), Inches(0.8)
    )
    tf2 = sub_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = 'Fiscal Year 2025 — Global Operations'
    run2.font.name = 'Calibri'
    run2.font.size = Pt(24)
    run2.font.color.rgb = RGBColor(0xCC, 0xD6, 0xE8)

    # Date text box
    date_box = slide1.shapes.add_textbox(
        Inches(4.5), Inches(5.0),
        Inches(4.5), Inches(0.5)
    )
    tf3 = date_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = 'December 10, 2025'
    run3.font.name = 'Calibri'
    run3.font.size = Pt(16)
    run3.font.color.rgb = RGBColor(0xFF, 0xA5, 0x00)

    # ─── Slide 2: Agenda ───
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])

    # Light background
    slide2_bg = slide2.shapes.add_shape(
        1, Inches(0), Inches(0),
        Inches(13.33), Inches(7.5)
    )
    slide2_bg.fill.solid()
    slide2_bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
    slide2_bg.line.fill.background()

    # Section header bar
    header_bar = slide2.shapes.add_shape(
        1, Inches(0), Inches(0),
        Inches(13.33), Inches(1.1)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x4A)
    header_bar.line.fill.background()

    header_text = slide2.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(6.0), Inches(0.8)
    )
    p_h = header_text.text_frame.paragraphs[0]
    run_h = p_h.add_run()
    run_h.text = 'Agenda'
    run_h.font.name = 'Calibri'
    run_h.font.size = Pt(32)
    run_h.font.bold = True
    run_h.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Agenda items
    agenda_items = [
        '1. Financial Performance Overview',
        '2. Regional Sales Analysis',
        '3. Product Portfolio Update',
        '4. Headcount & Talent Strategy',
        '5. Key Risks and Mitigation Plans',
        '6. FY2026 Strategic Priorities',
    ]
    agenda_box = slide2.shapes.add_textbox(
        Inches(1.2), Inches(1.5),
        Inches(10.0), Inches(5.5)
    )
    tf_agenda = agenda_box.text_frame
    tf_agenda.word_wrap = True
    for i, item in enumerate(agenda_items):
        if i == 0:
            para = tf_agenda.paragraphs[0]
        else:
            para = tf_agenda.add_paragraph()
        run_a = para.add_run()
        run_a.text = item
        run_a.font.name = 'Calibri'
        run_a.font.size = Pt(20)
        run_a.font.color.rgb = RGBColor(0x1A, 0x2A, 0x4A)
        para.space_before = Pt(8)

    # ─── Slide 3: Financial Highlights ───
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])

    slide3_bg = slide3.shapes.add_shape(
        1, Inches(0), Inches(0),
        Inches(13.33), Inches(7.5)
    )
    slide3_bg.fill.solid()
    slide3_bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
    slide3_bg.line.fill.background()

    header_bar3 = slide3.shapes.add_shape(
        1, Inches(0), Inches(0),
        Inches(13.33), Inches(1.1)
    )
    header_bar3.fill.solid()
    header_bar3.fill.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x4A)
    header_bar3.line.fill.background()

    header_text3 = slide3.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(9.0), Inches(0.8)
    )
    p_h3 = header_text3.text_frame.paragraphs[0]
    run_h3 = p_h3.add_run()
    run_h3.text = 'Financial Performance Overview'
    run_h3.font.name = 'Calibri'
    run_h3.font.size = Pt(28)
    run_h3.font.bold = True
    run_h3.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Add a summary table
    table_shape = slide3.shapes.add_table(
        5, 4,
        Inches(0.8), Inches(1.4),
        Inches(11.5), Inches(3.5)
    )
    table = table_shape.table
    headers = ['Region', 'Revenue (M)', 'Growth YoY', 'Target Achieved']
    data = [
        ['North America', '$142.8', '+12.4%', '104%'],
        ['Europe', '$89.3', '+8.7%', '97%'],
        ['Asia Pacific', '$67.5', '+21.3%', '112%'],
        ['Latin America', '$31.2', '+5.1%', '88%'],
    ]
    for col, h in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.size = Pt(14)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x4A)

    for row_idx, row_data in enumerate(data):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(13)
                run.font.color.rgb = RGBColor(0x1A, 0x2A, 0x4A)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xED, 0xF5)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Footer note
    footer_box = slide3.shapes.add_textbox(
        Inches(0.8), Inches(5.2),
        Inches(11.5), Inches(0.5)
    )
    p_footer = footer_box.text_frame.paragraphs[0]
    run_footer = p_footer.add_run()
    run_footer.text = 'Note: All figures are unaudited estimates as of November 30, 2025.'
    run_footer.font.name = 'Calibri'
    run_footer.font.size = Pt(11)
    run_footer.font.italic = True
    run_footer.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

    # Save
    os.makedirs(os.path.dirname(DESKTOP_OUTPUT), exist_ok=True)
    prs.save(OUTPUT)
    shutil.copy(OUTPUT, DESKTOP_OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Also copied to: {DESKTOP_OUTPUT}')
    print(f'Slide 1 background rectangle: UNLOCKED (ready for task)')


create_initial()
