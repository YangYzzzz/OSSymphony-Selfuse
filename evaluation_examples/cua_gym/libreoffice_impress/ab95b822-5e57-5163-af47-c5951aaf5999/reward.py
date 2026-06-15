"""
FINAL REWARD SCRIPT - SUCCESS
Task: In LibreOffice Impress, I’m working on slide 282 and want to push “Text Box 1” all the way to the right edge, then give it a precise 1.0 cm right-side margin. What steps get that done?
Generated: 2025-09-10 21:28:58
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation

"""
Reward Script for LibreOffice Impress Task Verification
Task: On slide 282, “Text Box 1” must be pushed to the right edge and given an exact 1.0 cm right-side margin.
This script checks:
 1. The presentation exists and loads correctly (no points – prerequisite)
 2. Slide 282 exists (no points – prerequisite)
 3. A shape named exactly “Text Box 1” is present on slide 282 (0.2 pts)
 4. The shape’s right edge is flush with the slide’s right edge (±0.22 cm tolerance) (0.4 pts)
 5. The shape’s text-frame right margin is 1.0 cm (±0.11 cm tolerance) (0.4 pts)
Score is progressive and capped at 1.0.
"""

def verify_impress_task(file_path: str) -> float:
    print(f"Verifying task for file: {file_path}")
    total_score = 0.0
    MAX_SCORE = 1.0

    # --- 1. Load presentation (prerequisite, zero points) -------------------
    if not os.path.exists(file_path):
        print("✗ File not found")
        return 0.0
    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded PPTX with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Unable to load presentation: {e}")
        return 0.0

    # --- 2. Ensure slide 282 exists (prerequisite) --------------------------
    target_index = 281  # zero-based index for slide 282
    if len(prs.slides) <= target_index:
        print("✗ Slide 282 not found in presentation")
        return 0.0
    slide = prs.slides[target_index]
    slide_width = prs.slide_width
    print(f"✓ Accessed slide 282  (slide width = {slide_width} EMU)")

    # --- 3. Locate shape named “Text Box 1” ---------------------------------
    target_shape = None
    for shape in slide.shapes:
        if shape.name.strip().lower() == "text box 1":
            target_shape = shape
            break
    if target_shape is None:
        print("✗ Shape named ‘Text Box 1’ not found on slide 282")
        return 0.0
    print(f"✓ Found target shape: {target_shape.name}")
    total_score += 0.2  # Identification success

    # --- 4. Verify right-edge alignment -------------------------------------
    right_edge_gap = slide_width - (target_shape.left + target_shape.width)
    print(f"Right-edge gap: {right_edge_gap} EMU")
    ALIGN_TOL = 20000  # ≈0.22 cm tolerance
    if abs(right_edge_gap) <= ALIGN_TOL:
        print("✓ Shape aligned to right edge (within tolerance)")
        total_score += 0.4
    else:
        print("✗ Shape is not correctly aligned to the right edge")

    # --- 5. Verify 1.0 cm right margin --------------------------------------
    if not target_shape.has_text_frame:
        print("✗ Target shape lacks a text frame – cannot verify margin")
    else:
        margin_right = target_shape.text_frame.margin_right
        print(f"Right margin: {margin_right} EMU")
        EXPECTED_MARGIN = 360000   # 1 cm in EMU
        MARGIN_TOL = 40000         # ≈0.11 cm tolerance
        if abs(margin_right - EXPECTED_MARGIN) <= MARGIN_TOL:
            print("✓ Right margin is 1.0 cm (within tolerance)")
            total_score += 0.4
        else:
            print("✗ Right margin is not 1.0 cm")

    # -----------------------------------------------------------------------
    final_score = min(total_score, MAX_SCORE)
    print(f"Total score: {final_score}/{MAX_SCORE}")
    print(f"REWARD: {final_score}")
    return final_score


# -------------------- Script Execution Entry Point -------------------------
if __name__ == "__main__":
    verify_impress_task("/home/user/in_libreoffice_impress_im_working_on_slide_282_and_want_to_push_text_box_1_all_the_way_to_the_right__golden.pptx")
