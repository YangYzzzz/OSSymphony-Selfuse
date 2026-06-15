"""
Initial Setup: Create 8-slide sales presentation with slide 5 titled 'Choose Your Plan' (empty content area)
Task ID: impress_sales_043
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_043'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                alignment=PP_ALIGN.LEFT, color=None, font_name="Calibri"):
    """Helper to add a formatted text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullet_paragraph(text_frame, text, font_size=14, color=None, font_name="Calibri", level=0):
    """Add a bullet paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    if color:
        run.font.color.rgb = color
    return p


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_textbox(slide1, Inches(1.5), Inches(1.5), Inches(10), Inches(2),
                "CloudSync Pro", font_size=44, bold=True,
                alignment=PP_ALIGN.CENTER, color=RGBColor(0x1B, 0x3A, 0x5C))
    add_textbox(slide1, Inches(2.5), Inches(3.8), Inches(8), Inches(1.2),
                "Enterprise Cloud Management Platform", font_size=24,
                alignment=PP_ALIGN.CENTER, color=RGBColor(0x55, 0x55, 0x55))
    add_textbox(slide1, Inches(3.5), Inches(5.5), Inches(6), Inches(0.8),
                "Q2 2025 Sales Deck  |  Confidential", font_size=14,
                alignment=PP_ALIGN.CENTER, color=RGBColor(0x99, 0x99, 0x99))

    # ---- Slide 2: Market Opportunity ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide2, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
                "Market Opportunity", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))
    txBox2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = "The global cloud management market is projected to reach $47.3B by 2027"
    tf2.paragraphs[0].runs[0].font.size = Pt(16)
    tf2.paragraphs[0].runs[0].font.name = "Calibri"
    for bullet_text in [
        "78% of enterprises plan to increase cloud spending in the next 12 months",
        "Average company uses 3.4 cloud providers, creating management complexity",
        "Security and compliance remain the #1 concern for cloud adoption (Gartner 2025)",
        "Multi-cloud orchestration tools market growing at 24.6% CAGR",
    ]:
        add_bullet_paragraph(tf2, bullet_text, font_size=16)

    # ---- Slide 3: Product Overview ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide3, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
                "Product Overview", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))
    features = [
        ("Unified Dashboard", "Single pane of glass across AWS, Azure, and GCP"),
        ("Smart Autoscaling", "AI-driven resource optimization reducing costs by up to 35%"),
        ("Compliance Engine", "Automated auditing for SOC 2, HIPAA, and GDPR frameworks"),
        ("Disaster Recovery", "One-click failover with 99.99% uptime SLA"),
    ]
    y_pos = Inches(1.8)
    for title, desc in features:
        add_textbox(slide3, Inches(1.2), y_pos, Inches(10), Inches(0.5),
                    title, font_size=20, bold=True, color=RGBColor(0x2E, 0x75, 0xB6))
        add_textbox(slide3, Inches(1.2), y_pos + Inches(0.5), Inches(10), Inches(0.5),
                    desc, font_size=15, color=RGBColor(0x44, 0x44, 0x44))
        y_pos += Inches(1.2)

    # ---- Slide 4: Customer Success ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide4, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
                "Customer Success Stories", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))
    testimonials = [
        ('"CloudSync Pro reduced our infrastructure costs by 42% in the first quarter."',
         "- Rachel Torres, VP Engineering, DataFlow Inc."),
        ('"The compliance engine saved our team 200+ hours during our SOC 2 audit."',
         "- James Park, CISO, MedSecure Health"),
        ('"We consolidated 5 monitoring tools into one. The ROI was immediate."',
         "- Priya Sharma, CTO, NexGen Logistics"),
    ]
    y_pos = Inches(1.8)
    for quote, attribution in testimonials:
        add_textbox(slide4, Inches(1.2), y_pos, Inches(10), Inches(0.6),
                    quote, font_size=16, color=RGBColor(0x33, 0x33, 0x33))
        add_textbox(slide4, Inches(1.5), y_pos + Inches(0.6), Inches(9), Inches(0.4),
                    attribution, font_size=13, color=RGBColor(0x88, 0x88, 0x88))
        y_pos += Inches(1.3)

    # ---- Slide 5: Choose Your Plan (EMPTY content area - task target) ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide5, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
                "Choose Your Plan", font_size=32, bold=True,
                alignment=PP_ALIGN.CENTER, color=RGBColor(0x1B, 0x3A, 0x5C))
    # Content area is intentionally empty - the agent must add the 3-column pricing layout

    # ---- Slide 6: Implementation Timeline ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide6, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
                "Implementation Timeline", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))
    phases = [
        ("Week 1-2", "Discovery & Architecture Review"),
        ("Week 3-4", "Platform Deployment & Configuration"),
        ("Week 5-6", "Data Migration & Integration"),
        ("Week 7-8", "User Training & Go-Live Support"),
    ]
    y_pos = Inches(1.8)
    for phase_label, phase_desc in phases:
        add_textbox(slide6, Inches(1.2), y_pos, Inches(2.5), Inches(0.5),
                    phase_label, font_size=18, bold=True, color=RGBColor(0xFF, 0x6B, 0x35))
        add_textbox(slide6, Inches(4.0), y_pos, Inches(7), Inches(0.5),
                    phase_desc, font_size=18, color=RGBColor(0x44, 0x44, 0x44))
        y_pos += Inches(0.8)

    # ---- Slide 7: ROI Calculator ----
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide7, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
                "Expected ROI", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))
    metrics = [
        ("35%", "Average cost reduction"),
        ("200+", "Hours saved per audit cycle"),
        ("99.99%", "Uptime SLA guarantee"),
        ("< 6 mo", "Typical payback period"),
    ]
    x_pos = Inches(0.8)
    for value, label in metrics:
        add_textbox(slide7, x_pos, Inches(2.5), Inches(2.8), Inches(1),
                    value, font_size=36, bold=True,
                    alignment=PP_ALIGN.CENTER, color=RGBColor(0x2E, 0x75, 0xB6))
        add_textbox(slide7, x_pos, Inches(3.8), Inches(2.8), Inches(0.6),
                    label, font_size=14,
                    alignment=PP_ALIGN.CENTER, color=RGBColor(0x66, 0x66, 0x66))
        x_pos += Inches(3.0)

    # ---- Slide 8: Contact & Next Steps ----
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide8, Inches(1.5), Inches(1.5), Inches(10), Inches(1.5),
                "Ready to Get Started?", font_size=40, bold=True,
                alignment=PP_ALIGN.CENTER, color=RGBColor(0x1B, 0x3A, 0x5C))
    add_textbox(slide8, Inches(2.5), Inches(3.5), Inches(8), Inches(0.8),
                "Contact our sales team for a personalized demo", font_size=20,
                alignment=PP_ALIGN.CENTER, color=RGBColor(0x55, 0x55, 0x55))
    add_textbox(slide8, Inches(3), Inches(4.8), Inches(7), Inches(0.5),
                "sales@cloudsyncpro.com  |  1-800-CLOUD-55", font_size=16,
                alignment=PP_ALIGN.CENTER, color=RGBColor(0x2E, 0x75, 0xB6))
    add_textbox(slide8, Inches(3), Inches(5.6), Inches(7), Inches(0.5),
                "www.cloudsyncpro.com/demo", font_size=14,
                alignment=PP_ALIGN.CENTER, color=RGBColor(0x88, 0x88, 0x88))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
