"""
Reward Script: Insert 3-column pricing layout on slide 5
Task ID: impress_sales_043
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Three column text boxes with correct headers
  Component 2 (0.25): Correct prices in each column
  Component 3 (0.25): Each column has 4 feature bullets
  Component 4 (0.25): Professional column has #FF6B35 ~3pt border
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_043'


def find_column_shapes(slide):
    """Find the three column text boxes on slide 5 (excluding the title).
    Returns dict mapping header text -> shape, or empty dict if not found.
    """
    columns = {}
    expected_headers = {'Starter', 'Professional', 'Enterprise'}
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        paras = shape.text_frame.paragraphs
        # Get first non-empty paragraph text
        first_text = None
        for p in paras:
            t = p.text.strip()
            if t:
                first_text = t
                break
        if first_text in expected_headers:
            columns[first_text] = shape
    return columns


def get_shape_border(shape):
    """Return (color_hex, width_emu) of the shape border, or (None, None) if no border."""
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    sp = shape._element
    ln = sp.find('.//a:ln', nsmap)
    if ln is None:
        return None, None
    width = ln.get('w')
    width = int(width) if width else None
    fill = ln.find('a:solidFill', nsmap)
    if fill is None:
        return None, width
    srgb = fill.find('a:srgbClr', nsmap)
    if srgb is not None:
        return srgb.get('val'), width
    return None, width


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # slide 5 (0-indexed)

    # Find the three column shapes
    columns = find_column_shapes(slide)

    # Component 1: Three column text boxes with correct headers (0.25 points)
    try:
        found_headers = set(columns.keys())
        expected = {'Starter', 'Professional', 'Enterprise'}
        if found_headers == expected:
            print(f"PASS: Component 1 — All three column headers found: {found_headers} (0.25 pts)")
            total_score += 0.25
        else:
            missing = expected - found_headers
            print(f"FAIL: Component 1 — Missing headers: {missing}. Found: {found_headers}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Correct prices in each column (0.25 points)
    try:
        price_map = {
            'Starter': '$29/mo',
            'Professional': '$79/mo',
            'Enterprise': '$199/mo',
        }
        prices_correct = 0
        for name, expected_price in price_map.items():
            if name not in columns:
                print(f"FAIL: Component 2 — Column '{name}' not found, cannot check price")
                continue
            shape = columns[name]
            # Price is typically in the second non-empty paragraph
            found_price = False
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text == expected_price:
                    found_price = True
                    break
            if found_price:
                prices_correct += 1
                print(f"  Price OK: {name} = {expected_price}")
            else:
                # Collect all para texts for debug
                all_texts = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
                print(f"  Price FAIL: {name} expected '{expected_price}', found paragraphs: {all_texts}")

        if prices_correct == 3:
            print(f"PASS: Component 2 — All three prices correct (0.25 pts)")
            total_score += 0.25
        elif prices_correct > 0:
            partial = round(0.25 * prices_correct / 3, 3)
            print(f"PARTIAL: Component 2 — {prices_correct}/3 prices correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No prices correct")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Each column has 4 feature bullet lines (0.25 points)
    try:
        bullets_correct = 0
        for name in ['Starter', 'Professional', 'Enterprise']:
            if name not in columns:
                print(f"FAIL: Component 3 — Column '{name}' not found")
                continue
            shape = columns[name]
            # Feature bullets: non-empty paragraphs that are NOT the header and NOT the price
            paras = shape.text_frame.paragraphs
            # Identify feature lines: skip header (first non-empty), price (contains $/mo),
            # and count remaining non-empty paragraphs
            feature_lines = []
            header_found = False
            price_found = False
            for p in paras:
                text = p.text.strip()
                if not text:
                    continue
                if not header_found and text in ['Starter', 'Professional', 'Enterprise']:
                    header_found = True
                    continue
                if not price_found and '/mo' in text:
                    price_found = True
                    continue
                feature_lines.append(text)

            if len(feature_lines) >= 4:
                bullets_correct += 1
                print(f"  Bullets OK: {name} has {len(feature_lines)} features")
            else:
                print(f"  Bullets FAIL: {name} has {len(feature_lines)} features, expected 4")

        if bullets_correct == 3:
            print(f"PASS: Component 3 — All columns have 4+ feature bullets (0.25 pts)")
            total_score += 0.25
        elif bullets_correct > 0:
            partial = round(0.25 * bullets_correct / 3, 3)
            print(f"PARTIAL: Component 3 — {bullets_correct}/3 columns have correct bullets ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No columns have correct bullet count")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Professional column has #FF6B35 border ~3pt width (0.25 points)
    try:
        if 'Professional' not in columns:
            print(f"FAIL: Component 4 — Professional column not found")
        else:
            prof_shape = columns['Professional']
            color, width = get_shape_border(prof_shape)

            border_color_ok = False
            border_width_ok = False

            if color is not None:
                # Check color matches FF6B35 (case-insensitive)
                if color.upper() == 'FF6B35':
                    border_color_ok = True
                    print(f"  Border color OK: #{color}")
                else:
                    print(f"  Border color FAIL: expected #FF6B35, found #{color}")
            else:
                print(f"  Border color FAIL: no border color found")

            if width is not None:
                # 3pt = 38100 EMU (1pt = 12700 EMU). Allow 20% tolerance.
                expected_w = 38100  # 3pt
                if abs(width - expected_w) / expected_w <= 0.2:
                    border_width_ok = True
                    print(f"  Border width OK: {width} EMU (~{width/12700:.1f}pt)")
                else:
                    print(f"  Border width FAIL: expected ~38100 EMU (3pt), found {width} EMU (~{width/12700:.1f}pt)")
            else:
                print(f"  Border width FAIL: no border width found")

            # Also verify other columns do NOT have this border (to confirm it's specific)
            other_has_border = False
            for name in ['Starter', 'Enterprise']:
                if name in columns:
                    oc, ow = get_shape_border(columns[name])
                    if oc == 'FF6B35':
                        other_has_border = True
                        print(f"  WARNING: {name} also has #FF6B35 border")

            if border_color_ok and border_width_ok:
                print(f"PASS: Component 4 — Professional has #FF6B35 3pt border (0.25 pts)")
                total_score += 0.25
            elif border_color_ok:
                # Partial: color is right but width is off
                print(f"PARTIAL: Component 4 — Border color correct but width wrong (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 4 — Professional column missing expected border")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
