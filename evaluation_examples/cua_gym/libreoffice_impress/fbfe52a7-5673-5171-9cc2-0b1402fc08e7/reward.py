"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 189 I have a 4-column object named “Table 1.” Right now it’s sitting off to the left and the columns are uneven. In LibreOffice Impress, how do I (a) center that entire table on the slide—dead-center both horizontally and vertically—and (b) force every column in the table to the exact same width, e.g., 2.5 cm each?
Generated: 2025-09-10 18:15:14
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
from pptx import Presentation
from pptx.util import Cm


def verify_table_alignment_and_columns(file_path: str,
                                        table_name: str = "Table 1",
                                        slide_index: int = 188) -> float:
    """Verify that the specified table is centered on the slide and that
    all its columns are of equal width (2.5 cm).

    Parameters
    ----------
    file_path : str
        Path to the PPTX file to check.
    table_name : str, optional
        Name of the table shape to verify. Defaults to "Table 1".
    slide_index : int, optional
        Zero-based index of the slide that should contain the table.
        Defaults to 188 (i.e., slide 189 in 1-based terms).

    Returns
    -------
    float
        Reward score between 0.0 and 1.0.
    """

    print(f"Starting verification for file: {file_path}")
    max_score = 1.0
    score = 0.0

    # Point distribution (progressive scoring)
    presence_points = 0.25        # Table exists on correct slide
    horizontal_points = 0.25 / 2  # Horizontal centering
    vertical_points   = 0.25 / 2  # Vertical centering
    equal_width_points = 0.25 / 2 # All columns equal
    width_value_points = 0.25 / 2 # Columns exactly 2.5 cm

    # ------------------------------------------------------------------
    # 1. Load presentation ------------------------------------------------
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File does not exist")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded presentation with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # ------------------------------------------------------------------
    # 2. Validate expected slide index -----------------------------------
    # ------------------------------------------------------------------
    if slide_index >= len(prs.slides):
        print(f"✗ Expected slide index {slide_index} not present (total slides: {len(prs.slides)})")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[slide_index]
    print(f"✓ Accessed slide {slide_index + 1}")

    # ------------------------------------------------------------------
    # 3. Locate the target table by name ---------------------------------
    # ------------------------------------------------------------------
    target_shape = None
    for shp in slide.shapes:
        if shp.has_table and getattr(shp, "name", "") == table_name:
            target_shape = shp
            break

    if target_shape is None:
        print(f"✗ Table named '{table_name}' not found on slide {slide_index + 1}")
        print("REWARD: 0.0")
        return 0.0

    print(f"✓ Found table '{table_name}' on slide")
    score += presence_points

    # ------------------------------------------------------------------
    # 4. Check centering (horizontal & vertical) -------------------------
    # ------------------------------------------------------------------
    slide_w, slide_h = prs.slide_width, prs.slide_height
    shape_left, shape_top = target_shape.left, target_shape.top
    shape_w, shape_h = target_shape.width, target_shape.height

    # Desired positions for perfect centering
    desired_left = (slide_w - shape_w) / 2
    desired_top  = (slide_h - shape_h) / 2

    tolerance = 50_000  # ±50 000 EMU ≈ 1.4 mm tolerance

    # Horizontal centering check
    if abs(shape_left - desired_left) <= tolerance:
        print(f"✓ Table horizontally centered (left={shape_left}, expected≈{int(desired_left)})")
        score += horizontal_points
    else:
        print(f"✗ Table NOT horizontally centered (left={shape_left}, expected≈{int(desired_left)})")

    # Vertical centering check
    if abs(shape_top - desired_top) <= tolerance:
        print(f"✓ Table vertically centered   (top={shape_top},  expected≈{int(desired_top)})")
        score += vertical_points
    else:
        print(f"✗ Table NOT vertically centered (top={shape_top}, expected≈{int(desired_top)})")

    # ------------------------------------------------------------------
    # 5. Check column widths --------------------------------------------
    # ------------------------------------------------------------------
    tbl = target_shape.table
    col_widths = [col.width for col in tbl.columns]
    print(f"Column widths (EMU): {col_widths}")

    # 5a. All columns equal width?
    if len(set(col_widths)) == 1:
        print("✓ All columns have equal width")
        score += equal_width_points
    else:
        print("✗ Columns are NOT all equal width")

    # 5b. Width equals 2.5 cm?
    expected_width = Cm(2.5).emu  # 2.5 cm in EMUs
    if all(abs(w - expected_width) <= tolerance for w in col_widths):
        print(f"✓ Each column width ≈ 2.5 cm (≈{expected_width} EMU)")
        score += width_value_points
    else:
        print(f"✗ Column widths do NOT match 2.5 cm (≈{expected_width} EMU)")

    # ------------------------------------------------------------------
    # 6. Final scoring ----------------------------------------------------
    # ------------------------------------------------------------------
    final_score = min(score, max_score)
    print(f"Final score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score


# ----------------------------------------------------------------------
# Execute verification when run as a script
# ----------------------------------------------------------------------
if __name__ == "__main__":
    verify_table_alignment_and_columns(
        "/home/user/on_slide_189_i_have_a_4_column_object_named_table_1_right_now_its_sitting_off_to_the_left_and_the_co_golden.pptx"
    )
