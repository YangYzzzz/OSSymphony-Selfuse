"""
Initial Setup: 10-slide product launch deck with white backgrounds
Task ID: osworld_impress_all_slides_background_008
Domain: libreoffice_impress

Creates a product launch presentation with 10 slides, all with white backgrounds
and dark text. The agent task is to apply a solid #333333 background to all slides.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_all_slides_background_008'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, text, left, top, width, height,
                 font_size=24, bold=False, color=RGBColor(0x33, 0x33, 0x33),
                 alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def set_white_background(slide):
    """Set a solid white background on the slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    dark = RGBColor(0x33, 0x33, 0x33)
    accent = RGBColor(0x1A, 0x78, 0xC2)
    gray = RGBColor(0x55, 0x55, 0x55)

    # -------------------------------------------------------------------
    # Slide 1: Title Slide
    # -------------------------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_white_background(slide1)
    add_text_box(slide1, "NovaSpark X1", Inches(1.5), Inches(2.0),
                 Inches(10), Inches(1.2), font_size=44, bold=True,
                 color=dark, alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, "Product Launch 2025", Inches(1.5), Inches(3.4),
                 Inches(10), Inches(0.8), font_size=26, bold=False,
                 color=accent, alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, "Transforming the way teams collaborate",
                 Inches(2.0), Inches(4.4), Inches(9), Inches(0.6),
                 font_size=18, color=gray, alignment=PP_ALIGN.CENTER)

    # -------------------------------------------------------------------
    # Slide 2: Agenda
    # -------------------------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(slide2)
    add_text_box(slide2, "Agenda", Inches(1.0), Inches(0.5),
                 Inches(11), Inches(0.8), font_size=32, bold=True, color=dark)
    items = [
        "1.  Company Overview",
        "2.  Problem Statement",
        "3.  Product Vision",
        "4.  Key Features",
        "5.  Market Opportunity",
        "6.  Competitive Landscape",
        "7.  Roadmap",
        "8.  Team",
        "9.  Financials",
        "10. Call to Action",
    ]
    for i, item in enumerate(items):
        add_text_box(slide2, item, Inches(1.5), Inches(1.4 + i * 0.52),
                     Inches(10), Inches(0.5), font_size=16, color=gray)

    # -------------------------------------------------------------------
    # Slide 3: Company Overview
    # -------------------------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(slide3)
    add_text_box(slide3, "Company Overview", Inches(1.0), Inches(0.4),
                 Inches(11), Inches(0.8), font_size=32, bold=True, color=dark)
    add_text_box(slide3,
                 "Founded in 2019, NovaSpark Labs is a San Francisco-based software company "
                 "specializing in AI-powered collaboration tools for enterprise teams.",
                 Inches(1.0), Inches(1.4), Inches(11), Inches(1.2),
                 font_size=18, color=gray)
    stats = [
        ("250+", "Enterprise Clients"),
        ("$42M", "ARR (2024)"),
        ("180", "Employees"),
        ("12", "Countries"),
    ]
    for i, (val, label) in enumerate(stats):
        x = Inches(1.0 + i * 2.8)
        add_text_box(slide3, val, x, Inches(3.2), Inches(2.4), Inches(0.7),
                     font_size=28, bold=True, color=accent, alignment=PP_ALIGN.CENTER)
        add_text_box(slide3, label, x, Inches(3.95), Inches(2.4), Inches(0.5),
                     font_size=14, color=gray, alignment=PP_ALIGN.CENTER)

    # -------------------------------------------------------------------
    # Slide 4: Problem Statement
    # -------------------------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(slide4)
    add_text_box(slide4, "The Problem", Inches(1.0), Inches(0.4),
                 Inches(11), Inches(0.8), font_size=32, bold=True, color=dark)
    problems = [
        "Teams spend 32% of their workday switching between disconnected tools",
        "Average enterprise uses 14+ separate SaaS applications daily",
        "Critical context is lost across email, chat, and project management silos",
        "Employee onboarding takes 3× longer due to fragmented knowledge bases",
    ]
    for i, prob in enumerate(problems):
        add_text_box(slide4, f"•  {prob}", Inches(1.2), Inches(1.5 + i * 1.1),
                     Inches(11), Inches(0.9), font_size=17, color=gray)

    # -------------------------------------------------------------------
    # Slide 5: Product Vision
    # -------------------------------------------------------------------
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(slide5)
    add_text_box(slide5, "Our Vision", Inches(1.0), Inches(0.4),
                 Inches(11), Inches(0.8), font_size=32, bold=True, color=dark)
    add_text_box(slide5,
                 "\"A unified intelligent workspace where every team member has instant "
                 "access to the right information, the right tools, and the right people — "
                 "exactly when they need them.\"",
                 Inches(1.5), Inches(1.6), Inches(10), Inches(2.0),
                 font_size=20, color=dark, alignment=PP_ALIGN.CENTER)
    add_text_box(slide5,
                 "NovaSpark X1 brings together project management, real-time communication, "
                 "document collaboration, and AI assistance into a single coherent platform.",
                 Inches(1.0), Inches(4.0), Inches(11), Inches(1.4),
                 font_size=17, color=gray)

    # -------------------------------------------------------------------
    # Slide 6: Key Features
    # -------------------------------------------------------------------
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(slide6)
    add_text_box(slide6, "Key Features", Inches(1.0), Inches(0.4),
                 Inches(11), Inches(0.8), font_size=32, bold=True, color=dark)
    features = [
        ("AI Copilot", "Contextual suggestions, meeting summaries, and smart drafts powered by GPT-4o"),
        ("Live Canvas", "Real-time collaborative whiteboard with 200+ templates and diagramming tools"),
        ("Smart Docs", "Version-controlled documents with inline task assignments and approvals"),
        ("Unified Inbox", "Single view for messages, notifications, and action items across all integrations"),
        ("Analytics Hub", "Team productivity dashboards with burndown charts and cycle time tracking"),
    ]
    for i, (title, desc) in enumerate(features):
        add_text_box(slide6, title, Inches(1.2), Inches(1.3 + i * 1.05),
                     Inches(3.0), Inches(0.45), font_size=16, bold=True, color=accent)
        add_text_box(slide6, desc, Inches(4.3), Inches(1.3 + i * 1.05),
                     Inches(8.0), Inches(0.45), font_size=15, color=gray)

    # -------------------------------------------------------------------
    # Slide 7: Market Opportunity
    # -------------------------------------------------------------------
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(slide7)
    add_text_box(slide7, "Market Opportunity", Inches(1.0), Inches(0.4),
                 Inches(11), Inches(0.8), font_size=32, bold=True, color=dark)
    market_data = [
        ("$58B", "Total Addressable Market (2025)"),
        ("$14B", "Serviceable Addressable Market"),
        ("19% CAGR", "Projected growth through 2030"),
        ("82%", "Enterprise adoption rate of collaboration tools"),
    ]
    for i, (val, label) in enumerate(market_data):
        x = Inches(0.6 + (i % 2) * 6.2)
        y = Inches(1.8 + (i // 2) * 2.2)
        add_text_box(slide7, val, x, y, Inches(5.5), Inches(0.8),
                     font_size=30, bold=True, color=accent, alignment=PP_ALIGN.CENTER)
        add_text_box(slide7, label, x, y + Inches(0.85), Inches(5.5), Inches(0.5),
                     font_size=15, color=gray, alignment=PP_ALIGN.CENTER)

    # -------------------------------------------------------------------
    # Slide 8: Roadmap
    # -------------------------------------------------------------------
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(slide8)
    add_text_box(slide8, "Product Roadmap", Inches(1.0), Inches(0.4),
                 Inches(11), Inches(0.8), font_size=32, bold=True, color=dark)
    roadmap = [
        ("Q1 2025", "GA Launch", "Core platform, AI Copilot v1, Slack/Jira integrations"),
        ("Q2 2025", "Expansion", "Mobile apps, advanced analytics, Salesforce connector"),
        ("Q3 2025", "Enterprise+", "SSO, SOC 2 Type II, custom workflows, API v2"),
        ("Q4 2025", "AI Platform", "AI agents, automated workflows, marketplace launch"),
    ]
    for i, (quarter, milestone, detail) in enumerate(roadmap):
        y = Inches(1.4 + i * 1.35)
        add_text_box(slide8, quarter, Inches(0.8), y, Inches(1.4), Inches(0.4),
                     font_size=14, bold=True, color=accent)
        add_text_box(slide8, milestone, Inches(2.4), y, Inches(2.5), Inches(0.4),
                     font_size=15, bold=True, color=dark)
        add_text_box(slide8, detail, Inches(5.1), y, Inches(7.5), Inches(0.4),
                     font_size=14, color=gray)

    # -------------------------------------------------------------------
    # Slide 9: Team
    # -------------------------------------------------------------------
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(slide9)
    add_text_box(slide9, "Our Team", Inches(1.0), Inches(0.4),
                 Inches(11), Inches(0.8), font_size=32, bold=True, color=dark)
    team = [
        ("Elena Vasquez", "CEO & Co-Founder", "Ex-Google PM, Stanford MBA"),
        ("David Okonkwo", "CTO & Co-Founder", "Ex-AWS Principal Engineer"),
        ("Priya Mehta", "VP Product", "Ex-Asana, 12 yrs in SaaS"),
        ("James Thornton", "VP Sales", "Ex-Salesforce, $200M quota"),
        ("Sofia Lindqvist", "Head of Design", "Ex-Figma, IDEO alum"),
        ("Marcus Chen", "VP Engineering", "Ex-Stripe, MIT CS"),
    ]
    for i, (name, role, bio) in enumerate(team):
        col = i % 3
        row = i // 3
        x = Inches(0.5 + col * 4.2)
        y = Inches(1.4 + row * 2.5)
        add_text_box(slide9, name, x, y, Inches(3.8), Inches(0.45),
                     font_size=16, bold=True, color=dark)
        add_text_box(slide9, role, x, y + Inches(0.48), Inches(3.8), Inches(0.4),
                     font_size=14, color=accent)
        add_text_box(slide9, bio, x, y + Inches(0.92), Inches(3.8), Inches(0.45),
                     font_size=13, color=gray)

    # -------------------------------------------------------------------
    # Slide 10: Call to Action
    # -------------------------------------------------------------------
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(slide10)
    add_text_box(slide10, "Join the Future of Work", Inches(1.0), Inches(1.8),
                 Inches(11), Inches(1.0), font_size=36, bold=True,
                 color=dark, alignment=PP_ALIGN.CENTER)
    add_text_box(slide10,
                 "Start your 30-day free trial today — no credit card required.",
                 Inches(1.5), Inches(3.1), Inches(10), Inches(0.7),
                 font_size=20, color=gray, alignment=PP_ALIGN.CENTER)
    add_text_box(slide10, "www.novasparklabs.com/x1",
                 Inches(2.0), Inches(4.0), Inches(9), Inches(0.6),
                 font_size=22, bold=True, color=accent, alignment=PP_ALIGN.CENTER)
    add_text_box(slide10, "contact@novasparklabs.com  |  +1 (415) 882-7700",
                 Inches(2.0), Inches(4.9), Inches(9), Inches(0.5),
                 font_size=16, color=gray, alignment=PP_ALIGN.CENTER)

    # -------------------------------------------------------------------
    # Save
    # -------------------------------------------------------------------
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Slide count: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
