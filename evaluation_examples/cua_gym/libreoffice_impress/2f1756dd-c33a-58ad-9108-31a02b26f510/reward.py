"""
Reward Script: Apply solid dark charcoal (#333333) background to ALL slides
Task ID: osworld_impress_all_slides_background_008
Domain: libreoffice_impress
Scoring:
  Component 1: Number of slides with solid #333333 background (partial credit, proportional, 0.7 pts)
  Component 2: All 10 slides have #333333 background (full completion bonus, 0.3 pts)
  Total: 1.0
"""

import os
from pptx import Presentation
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_all_slides_background_008'
TARGET_COLOR = '333333'
EXPECTED_SLIDE_COUNT = 10


def verify_task(file_path):
    """
    Verify that all 10 slides have a solid #333333 background applied.
    Initial state: all slides have #FFFFFF (white) backgrounds.
    Golden state: all slides have #333333 (dark charcoal) backgrounds.

    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition gate: verify expected slide count
    num_slides = len(prs.slides)
    if num_slides != EXPECTED_SLIDE_COUNT:
        print(f"PRECONDITION FAIL: Expected {EXPECTED_SLIDE_COUNT} slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Presentation has {num_slides} slides (expected {EXPECTED_SLIDE_COUNT})")

    # Component 1: Count slides with solid #333333 background (proportional partial credit, 0.7 pts)
    # This FAILS on initial (all slides have #FFFFFF) and PASSES on golden (all slides have #333333)
    slides_with_correct_bg = 0
    try:
        for i, slide in enumerate(prs.slides):
            fill = slide.background.fill
            fill_type = fill.type
            if fill_type == 1:  # solid fill
                try:
                    rgb = fill.fore_color.rgb
                    rgb_str = str(rgb).upper()
                    if rgb_str == TARGET_COLOR.upper():
                        slides_with_correct_bg += 1
                    else:
                        print(f"FAIL: Slide {i+1} has solid fill #{rgb_str}, expected #333333")
                except Exception as e:
                    print(f"FAIL: Slide {i+1} solid fill color could not be read: {e}")
            else:
                print(f"FAIL: Slide {i+1} has fill type {fill_type} (not solid), expected solid #333333")

        proportion = slides_with_correct_bg / EXPECTED_SLIDE_COUNT
        component1_score = round(0.7 * proportion, 4)
        print(f"Component 1: {slides_with_correct_bg}/{EXPECTED_SLIDE_COUNT} slides with #333333 background ({component1_score:.4f}/0.7 pts)")
        if slides_with_correct_bg > 0:
            total_score += component1_score

    except Exception as e:
        print(f"ERROR: Component 1 — background color check failed: {e}")

    # Component 2: All 10 slides have #333333 background (full completion check, 0.3 pts)
    # This FAILS on initial (0 slides have #333333) and PASSES on golden (all 10 have #333333)
    try:
        if slides_with_correct_bg == EXPECTED_SLIDE_COUNT:
            print(f"PASS: Component 2 — All {EXPECTED_SLIDE_COUNT} slides confirmed with #333333 background (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 — Only {slides_with_correct_bg}/{EXPECTED_SLIDE_COUNT} slides have #333333; all {EXPECTED_SLIDE_COUNT} required for full credit")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = round(min(total_score, 1.0), 4)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the given VM env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
