"""
Reward Script: Create 5-pointed star on slide 1 with gradient fill and border
Task ID: impress_ndo_040
Domain: libreoffice_impress
Scoring:
  - Component 1: Star shape exists on slide 1 (0.25 pts)
  - Component 2: Star is 7cm x 7cm and centered (0.25 pts)
  - Component 3: Gradient fill from #FFD700 to #FF8C00 (0.30 pts)
  - Component 4: 1pt black border (0.20 pts)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_040'

NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}


def find_star_shape_xml(pptx_path):
    """Find the 5-point star shape in slide 1 XML. Returns the sp element or None."""
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            with zf.open('ppt/slides/slide1.xml') as f:
                root = ET.parse(f).getroot()
    except Exception as e:
        print(f"ERROR: Cannot parse slide1.xml: {e}")
        return None

    for sp in root.findall('.//p:cSld/p:spTree/p:sp', NS):
        prstGeom = sp.find('.//a:prstGeom', NS)
        if prstGeom is not None:
            prst = prstGeom.get('prst', '')
            if prst == 'star5':
                return sp
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    if not os.path.exists(file_path):
        print(f"CRITICAL: File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    # Use python-pptx for basic shape checks, XML for fill/line details
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    if len(prs.slides) == 0:
        print("CRITICAL: No slides in presentation")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Find star via XML (more reliable for shape type detection)
    star_xml = find_star_shape_xml(file_path)

    # Component 1: 5-pointed star shape exists on slide 1 (0.25 points)
    try:
        if star_xml is not None:
            print("PASS: Component 1 — 5-point star shape (star5) found on slide 1 (0.25 pts)")
            total_score += 0.25
        else:
            print("FAIL: Component 1 — No 5-point star shape (star5) found on slide 1")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    if star_xml is None:
        # No star means all subsequent checks fail
        print("FAIL: Components 2-4 skipped (no star found)")
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: Star is 7cm x 7cm and centered on slide (0.25 points)
    try:
        xfrm = star_xml.find('.//a:xfrm', NS)
        if xfrm is None:
            print("FAIL: Component 2 — No transform found on star shape")
        else:
            ext = xfrm.find('a:ext', NS)
            off = xfrm.find('a:off', NS)

            cx = int(ext.get('cx', '0'))
            cy = int(ext.get('cy', '0'))
            x = int(off.get('x', '0'))
            y = int(off.get('y', '0'))

            # 7cm = 2520000 EMU
            TARGET_SIZE = 2520000
            SIZE_TOLERANCE = 0.05  # 5% tolerance

            width_ok = abs(cx - TARGET_SIZE) / TARGET_SIZE <= SIZE_TOLERANCE
            height_ok = abs(cy - TARGET_SIZE) / TARGET_SIZE <= SIZE_TOLERANCE

            # Centering: shape center should be at slide center
            shape_center_x = x + cx // 2
            shape_center_y = y + cy // 2
            slide_center_x = slide_width // 2
            slide_center_y = slide_height // 2

            # Allow 5% of slide dimension tolerance for centering
            center_x_ok = abs(shape_center_x - slide_center_x) <= slide_width * 0.05
            center_y_ok = abs(shape_center_y - slide_center_y) <= slide_height * 0.05

            size_ok = width_ok and height_ok
            centered = center_x_ok and center_y_ok

            if size_ok and centered:
                print(f"PASS: Component 2 — Star is {cx/360000:.1f}cm x {cy/360000:.1f}cm, centered (0.25 pts)")
                total_score += 0.25
            elif size_ok:
                print(f"PARTIAL: Component 2 — Size correct ({cx/360000:.1f}cm x {cy/360000:.1f}cm) but not centered "
                      f"(offset: {abs(shape_center_x - slide_center_x)}, {abs(shape_center_y - slide_center_y)}) (0.1 pts)")
                total_score += 0.1
            elif centered:
                print(f"PARTIAL: Component 2 — Centered but wrong size ({cx/360000:.1f}cm x {cy/360000:.1f}cm, expected 7.0cm) (0.1 pts)")
                total_score += 0.1
            else:
                print(f"FAIL: Component 2 — Size {cx/360000:.1f}cm x {cy/360000:.1f}cm (expected 7.0cm), not centered")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Gradient fill from #FFD700 to #FF8C00 (0.30 points)
    try:
        gradFill = star_xml.find('.//a:gradFill', NS)
        if gradFill is None:
            print("FAIL: Component 3 — No gradient fill on star shape")
        else:
            gs_list = gradFill.findall('.//a:gs', NS)
            if len(gs_list) < 2:
                print(f"FAIL: Component 3 — Gradient has {len(gs_list)} stops, expected at least 2")
            else:
                # Extract gradient stop colors
                stops = []
                for gs in gs_list:
                    pos = int(gs.get('pos', '0'))
                    clr_elem = gs.find('a:srgbClr', NS)
                    if clr_elem is not None:
                        color = clr_elem.get('val', '').upper()
                    else:
                        color = None
                    stops.append((pos, color))

                stops.sort(key=lambda x: x[0])
                print(f"  Gradient stops: {stops}")

                # Check first stop is gold (#FFD700) and last is dark orange (#FF8C00)
                first_color = stops[0][1]
                last_color = stops[-1][1]

                has_gold = first_color == 'FFD700'
                has_orange = last_color == 'FF8C00'

                if has_gold and has_orange:
                    print(f"PASS: Component 3 — Gradient from #FFD700 to #FF8C00 (0.30 pts)")
                    total_score += 0.30
                elif has_gold or has_orange:
                    matched = '#FFD700' if has_gold else '#FF8C00'
                    missing = '#FF8C00' if has_gold else '#FFD700'
                    print(f"PARTIAL: Component 3 — Has {matched} but missing {missing} (0.15 pts)")
                    total_score += 0.15
                else:
                    print(f"FAIL: Component 3 — Expected #FFD700 and #FF8C00, found {first_color} and {last_color}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: 1pt black (#000000) border (0.20 points)
    try:
        ln = star_xml.find('.//a:ln', NS)
        if ln is None:
            print("FAIL: Component 4 — No line/border on star shape")
        else:
            # Check line width: 1pt = 12700 EMU
            w = int(ln.get('w', '0'))
            TARGET_WIDTH = 12700
            width_ok = abs(w - TARGET_WIDTH) / TARGET_WIDTH <= 0.1  # 10% tolerance

            # Check line color
            solid_fill = ln.find('a:solidFill', NS)
            color_ok = False
            actual_color = None
            if solid_fill is not None:
                clr = solid_fill.find('a:srgbClr', NS)
                if clr is not None:
                    actual_color = clr.get('val', '').upper()
                    color_ok = actual_color == '000000'

            if width_ok and color_ok:
                print(f"PASS: Component 4 — 1pt ({w} EMU) black border (0.20 pts)")
                total_score += 0.20
            elif color_ok:
                print(f"PARTIAL: Component 4 — Black border but width {w} EMU (expected ~12700) (0.10 pts)")
                total_score += 0.10
            elif width_ok:
                print(f"PARTIAL: Component 4 — 1pt border but color {actual_color} (expected 000000) (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 4 — Width {w} EMU, color {actual_color} (expected 12700 EMU, 000000)")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
