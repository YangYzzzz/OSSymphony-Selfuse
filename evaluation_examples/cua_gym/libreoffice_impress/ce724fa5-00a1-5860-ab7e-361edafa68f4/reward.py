"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 251 is sticking out like a sore thumb—it’s still in the default title style. Could you flip that title over to Noto Serif, bump the size up to exactly 44 pt, and make it italic so it matches the rest of the deck?
Generated: 2025-09-10 19:17:38
Status: success
Model: azure-o3
Total Steps: 3
"""

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
import os


def verify_slide_251_title_format(pptx_path: str) -> float:
    """Verify that slide 251’s title is Noto Serif 44 pt italic.

    Returns a progressive score between 0.0-1.0 and prints detailed
    verification output in the required format ("REWARD: X.X").
    """

    score = 0.0
    max_score = 1.0

    # ------------------------------------------------------------------
    # 1. Load presentation ------------------------------------------------
    # ------------------------------------------------------------------
    if not os.path.exists(pptx_path):
        print(f"✗ File not found: {pptx_path}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # ------------------------------------------------------------------
    # 2. Ensure slide 251 exists (index 250) -----------------------------
    # ------------------------------------------------------------------
    slide_index = 250  # zero-based index
    if len(prs.slides) <= slide_index:
        print(f"✗ Slide 251 not found. Total slides: {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[slide_index]

    # ------------------------------------------------------------------
    # 3. Locate the title placeholder/shape -----------------------------
    # ------------------------------------------------------------------
    title_shape = None

    # Primary search: TITLE or CENTER_TITLE placeholders
    for shape in slide.shapes:
        if shape.has_text_frame and shape.is_placeholder and shape.placeholder_format.type in {
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
        }:
            title_shape = shape
            break

    # Fallback search: any shape whose name contains the word "title"
    if title_shape is None:
        for shape in slide.shapes:
            if shape.has_text_frame and "title" in shape.name.lower():
                title_shape = shape
                break

    if title_shape is None:
        print("✗ Title placeholder not found on slide 251")
        print("REWARD: 0.0")
        return 0.0

    print(f"✓ Found title shape: '{title_shape.name}'")

    # ------------------------------------------------------------------
    # 4. Inspect all text runs in the title -----------------------------
    # ------------------------------------------------------------------
    runs = [run for para in title_shape.text_frame.paragraphs for run in para.runs if run.text]
    if not runs:
        print("✗ Title has no text runs")
        print("REWARD: 0.0")
        return 0.0

    # Requirement checks -------------------------------------------------
    font_name_ok = all(run.font.name and "noto serif" in run.font.name.lower() for run in runs)
    size_ok = all(run.font.size and abs(run.font.size.pt - 44.0) < 0.1 for run in runs)
    italic_ok = all(run.font.italic is True for run in runs)

    # Scoring ------------------------------------------------------------
    if font_name_ok:
        print("✓ All title runs use 'Noto Serif' font")
        score += 0.30  # 30 %
    else:
        print("✗ Not all title runs use 'Noto Serif'")

    if size_ok:
        print("✓ All title runs are exactly 44 pt")
        score += 0.35  # 35 %
    else:
        print("✗ Title run sizes are not exactly 44 pt")

    if italic_ok:
        print("✓ All title runs are italic")
        score += 0.35  # 35 %
    else:
        print("✗ Title runs are not all italic")

    # Final score --------------------------------------------------------
    final_score = min(score, max_score)
    if abs(final_score - 1.0) < 1e-6:
        final_score = 1.0  # clean up floating-point artifacts
    final_score = round(final_score, 4)  # limit to 4 decimals for neatness

    print(f"Total score: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score


# ----------------------------------------------------------------------
# Execute verification when script is run directly ----------------------
# ----------------------------------------------------------------------
if __name__ == "__main__":
    verify_slide_251_title_format(
        "/home/user/slide_251_is_sticking_out_like_a_sore_thumbits_still_in_the_default_title_style_could_you_flip_that__golden.pptx"
    )

