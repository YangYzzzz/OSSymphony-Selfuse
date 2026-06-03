"""
Initial Setup: Create Onboarding_Flow.pptx with 6 slides. Slide 3 has only a title 'How It Works'.
Task ID: impress_design_078
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_design_078'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Onboarding Flow"
    slide1.placeholders[1].text = "A Comprehensive Guide for New Team Members"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    tb = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Agenda"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    items = [
        "Welcome & Introductions",
        "Company Overview",
        "How It Works",
        "Tools & Resources",
        "Q&A Session",
    ]
    for i, item in enumerate(items):
        p2 = tf.add_paragraph()
        p2.text = f"{i+1}. {item}"
        p2.space_before = Pt(6)
        r = p2.runs[0]
        r.font.size = Pt(18)
        r.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)

    # --- Slide 3: "How It Works" - title only, NO step content ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    tb3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "How It Works"
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.runs[0]
    run3.font.size = Pt(36)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # --- Slide 4: Tools & Resources ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    tb4 = slide4.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
    tf4 = tb4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Tools & Resources"
    p4.alignment = PP_ALIGN.LEFT
    run4 = p4.runs[0]
    run4.font.size = Pt(36)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    tools = [
        ("Slack", "Team communication and quick questions"),
        ("Confluence", "Documentation and knowledge base"),
        ("Jira", "Project tracking and task management"),
        ("GitHub", "Code repositories and version control"),
    ]
    for name, desc in tools:
        pt = tf4.add_paragraph()
        pt.space_before = Pt(8)
        rn = pt.add_run()
        rn.text = f"{name}: "
        rn.font.size = Pt(16)
        rn.font.bold = True
        rn.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)
        rd = pt.add_run()
        rd.text = desc
        rd.font.size = Pt(16)
        rd.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)

    # --- Slide 5: Timeline ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    tb5 = slide5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
    tf5 = tb5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Your First 90 Days"
    run5 = p5.runs[0]
    run5.font.size = Pt(36)
    run5.font.bold = True
    run5.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    milestones = [
        ("Week 1-2", "Complete orientation, meet your team, set up dev environment"),
        ("Week 3-4", "Shadow senior engineers, complete first code review"),
        ("Month 2", "Take ownership of a small feature, attend architecture meetings"),
        ("Month 3", "Lead a sprint task, present at team demo day"),
    ]
    for period, detail in milestones:
        pm = tf5.add_paragraph()
        pm.space_before = Pt(8)
        rp = pm.add_run()
        rp.text = f"{period}: "
        rp.font.size = Pt(16)
        rp.font.bold = True
        rp.font.color.rgb = RGBColor(0x2E, 0xCC, 0x71)
        rd2 = pm.add_run()
        rd2.text = detail
        rd2.font.size = Pt(16)
        rd2.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)

    # --- Slide 6: Q&A ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    tb6 = slide6.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf6 = tb6.text_frame
    p6 = tf6.paragraphs[0]
    p6.text = "Questions?"
    p6.alignment = PP_ALIGN.CENTER
    run6 = p6.runs[0]
    run6.font.size = Pt(48)
    run6.font.bold = True
    run6.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    p6b = tf6.add_paragraph()
    p6b.text = "Feel free to reach out to your onboarding buddy anytime!"
    p6b.alignment = PP_ALIGN.CENTER
    r6b = p6b.runs[0]
    r6b.font.size = Pt(20)
    r6b.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
