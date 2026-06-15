"""
Reward Script: Numbered step process with circles, titles, descriptions on slide 3
Task ID: impress_design_078
Domain: libreoffice_impress
Scoring:
  C1 (0.30) - Three circles with correct fill colors
  C2 (0.25) - Circle numbers 1/2/3 in 36pt bold white
  C3 (0.20) - Circle positions and sizes (x=1/5/9in, y=1.5in, 1.5in dia)
  C4 (0.15) - Three title textboxes below circles, 18pt bold
  C5 (0.10) - Three description textboxes below titles, 14pt
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_design_078'

# Expected circle colors (hex uppercase)
EXPECTED_COLORS = ['3498DB', '2ECC71', 'E74C3C']
# Expected circle x positions (EMU): 1in, 5in, 9in
EXPECTED_X_POS = [Inches(1), Inches(5), Inches(9)]
# Expected y position: 1.5in
EXPECTED_Y_POS = Inches(1.5)
# Expected diameter: 1.5in
EXPECTED_DIAMETER = Inches(1.5)

# Tolerance for position/size checks (5% relative)
def approx_equal(val1, val2, tolerance=0.05):
    if val1 == val2:
        return True
    if val1 == 0 and val2 == 0:
        return True
    if val1 == 0 or val2 == 0:
        return abs(val1 - val2) < Inches(0.2)
    return abs(val1 - val2) / max(abs(val1), abs(val2)) <= tolerance


def verify_task(file_path):
    """Verify task completion with progressive scoring. Returns float 0.0-1.0."""
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[2]  # Slide 3 (0-indexed)

    # Collect circles (ovals) on slide 3
    circles = []
    textboxes = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                from pptx.enum.shapes import MSO_SHAPE
                if shape.auto_shape_type == MSO_SHAPE.OVAL:
                    circles.append(shape)
            except Exception:
                # Fallback: check shape name or type value
                if hasattr(shape, 'auto_shape_type') and shape.auto_shape_type is not None:
                    if shape.auto_shape_type == 9:  # OVAL enum value
                        circles.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            # Collect non-title textboxes that are below the original position
            # The original "How It Works" textbox is at y~274320
            if shape.top > Inches(1):
                textboxes.append(shape)

    print(f"INFO: Found {len(circles)} circles and {len(textboxes)} new textboxes on slide 3")

    # Sort circles by x position (left to right)
    circles.sort(key=lambda s: s.left)
    # Sort textboxes by x then y
    textboxes.sort(key=lambda s: (s.left, s.top))

    # =========================================================
    # Component 1: Three circles with correct fill colors (0.30)
    # =========================================================
    try:
        if len(circles) >= 3:
            color_matches = 0
            for i, circle in enumerate(circles[:3]):
                try:
                    fill = circle.fill
                    if fill.type == 1:  # SOLID
                        actual_color = str(fill.fore_color.rgb).upper()
                        expected_color = EXPECTED_COLORS[i].upper()
                        if actual_color == expected_color:
                            color_matches += 1
                            print(f"PASS: Circle {i+1} fill color {actual_color} matches expected {expected_color}")
                        else:
                            print(f"FAIL: Circle {i+1} fill color {actual_color} != expected {expected_color}")
                    else:
                        print(f"FAIL: Circle {i+1} fill type is {fill.type}, not SOLID (1)")
                except Exception as e:
                    print(f"ERROR: Circle {i+1} fill check: {e}")

            if color_matches == 3:
                total_score += 0.30
                print(f"PASS: Component 1 -- All 3 circles have correct fill colors (0.30 pts)")
            elif color_matches > 0:
                partial = 0.10 * color_matches
                total_score += partial
                print(f"PARTIAL: Component 1 -- {color_matches}/3 circles have correct colors ({partial:.2f} pts)")
            else:
                print(f"FAIL: Component 1 -- No circles have correct fill colors")
        else:
            print(f"FAIL: Component 1 -- Found {len(circles)} circles, expected 3")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # =========================================================
    # Component 2: Numbers 1/2/3 in 36pt bold white (0.25)
    # =========================================================
    try:
        if len(circles) >= 3:
            number_matches = 0
            for i, circle in enumerate(circles[:3]):
                expected_num = str(i + 1)
                if circle.has_text_frame:
                    all_text = circle.text_frame.text.strip()
                    if all_text == expected_num:
                        # Check font properties on the matching run
                        run_checks = {'found': 0, 'bold': 0, 'size': 0, 'color': 0}
                        for para in circle.text_frame.paragraphs:
                            for run in para.runs:
                                if run.text.strip() == expected_num:
                                    run_checks['found'] = 1
                                    # Check bold
                                    if run.font.bold is True:
                                        run_checks['bold'] = 1
                                    else:
                                        print(f"FAIL: Circle {i+1} number not bold")
                                    # Check size (~36pt = 457200 EMU)
                                    if run.font.size is not None and approx_equal(run.font.size, Pt(36)):
                                        run_checks['size'] = 1
                                    else:
                                        print(f"FAIL: Circle {i+1} number size {run.font.size} != {Pt(36)}")
                                    # Check white color
                                    try:
                                        if run.font.color.type is not None:
                                            rgb = str(run.font.color.rgb).upper()
                                            if rgb == 'FFFFFF':
                                                run_checks['color'] = 1
                                            else:
                                                print(f"FAIL: Circle {i+1} number color {rgb} != FFFFFF")
                                    except Exception:
                                        pass
                        if sum(run_checks.values()) == 4:
                            number_matches += 1
                            print(f"PASS: Circle {i+1} has '{expected_num}' in 36pt bold white")
                    else:
                        print(f"FAIL: Circle {i+1} text '{all_text}' != expected '{expected_num}'")
                else:
                    print(f"FAIL: Circle {i+1} has no text frame")

            if number_matches == 3:
                total_score += 0.25
                print(f"PASS: Component 2 -- All 3 circles have correct numbers (0.25 pts)")
            elif number_matches > 0:
                partial = round(0.25 * number_matches / 3, 2)
                total_score += partial
                print(f"PARTIAL: Component 2 -- {number_matches}/3 correct ({partial:.2f} pts)")
            else:
                print(f"FAIL: Component 2 -- No circles have correct number text/formatting")
        else:
            print(f"FAIL: Component 2 -- Not enough circles")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # =========================================================
    # Component 3: Circle positions and sizes (0.20)
    # =========================================================
    try:
        if len(circles) >= 3:
            pos_matches = 0
            for i, circle in enumerate(circles[:3]):
                x_ok = approx_equal(circle.left, EXPECTED_X_POS[i])
                y_ok = approx_equal(circle.top, EXPECTED_Y_POS)
                w_ok = approx_equal(circle.width, EXPECTED_DIAMETER)
                h_ok = approx_equal(circle.height, EXPECTED_DIAMETER)

                if x_ok and y_ok and w_ok and h_ok:
                    pos_matches += 1
                    print(f"PASS: Circle {i+1} position/size correct")
                else:
                    print(f"FAIL: Circle {i+1} pos/size -- x:{x_ok}(actual={circle.left},exp={EXPECTED_X_POS[i]}) "
                          f"y:{y_ok}(actual={circle.top},exp={EXPECTED_Y_POS}) "
                          f"w:{w_ok}(actual={circle.width},exp={EXPECTED_DIAMETER}) "
                          f"h:{h_ok}(actual={circle.height},exp={EXPECTED_DIAMETER})")

            if pos_matches == 3:
                total_score += 0.20
                print(f"PASS: Component 3 -- All circles positioned correctly (0.20 pts)")
            elif pos_matches > 0:
                partial = round(0.20 * pos_matches / 3, 2)
                total_score += partial
                print(f"PARTIAL: Component 3 -- {pos_matches}/3 correct ({partial:.2f} pts)")
            else:
                print(f"FAIL: Component 3 -- No circles in correct position/size")
        else:
            print(f"FAIL: Component 3 -- Not enough circles")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # =========================================================
    # Component 4: Title textboxes below circles, 18pt bold (0.15)
    # =========================================================
    try:
        # Separate textboxes into titles (18pt bold, short text) vs descriptions (14pt, longer)
        # by checking font properties rather than just y position
        title_boxes = []
        desc_boxes_all = []
        for tb in textboxes:
            if not tb.has_text_frame:
                continue
            text = tb.text_frame.text.strip()
            if not text:
                continue
            # Check first non-empty run for font size
            for para in tb.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip() and run.font.size is not None:
                        if approx_equal(run.font.size, Pt(18)) and run.font.bold is True:
                            title_boxes.append(tb)
                        elif approx_equal(run.font.size, Pt(14)):
                            desc_boxes_all.append(tb)
                        break
                else:
                    continue
                break

        title_boxes.sort(key=lambda s: s.left)
        print(f"INFO: Identified {len(title_boxes)} title boxes and {len(desc_boxes_all)} desc boxes by font")

        title_count = 0
        if len(title_boxes) >= 3:
            for i, tb in enumerate(title_boxes[:3]):
                text = tb.text_frame.text.strip()
                if text:
                    title_count += 1
                    print(f"PASS: Title {i+1} '{text}' is 18pt bold")
                else:
                    print(f"FAIL: Title {i+1} is empty")
        else:
            print(f"FAIL: Found {len(title_boxes)} title boxes with 18pt bold, expected 3")

        if title_count == 3:
            total_score += 0.15
            print(f"PASS: Component 4 -- All 3 titles correct (0.15 pts)")
        elif title_count > 0:
            partial = round(0.15 * title_count / 3, 2)
            total_score += partial
            print(f"PARTIAL: Component 4 -- {title_count}/3 correct ({partial:.2f} pts)")
        else:
            print(f"FAIL: Component 4 -- No title textboxes found with correct formatting")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # =========================================================
    # Component 5: Description textboxes, 14pt (0.10)
    # =========================================================
    try:
        # Use desc_boxes_all identified by font in Component 4
        desc_boxes = sorted(desc_boxes_all, key=lambda s: s.left)

        desc_count = 0
        if len(desc_boxes) >= 3:
            for i, tb in enumerate(desc_boxes[:3]):
                if tb.has_text_frame:
                    text = tb.text_frame.text.strip()
                    if text and len(text) > 10:  # descriptions should be substantial
                        desc_count += 1
                        print(f"PASS: Description {i+1} text present, 14pt")
                    else:
                        print(f"FAIL: Description {i+1} text too short or empty: '{text[:40]}'")
        else:
            print(f"FAIL: Found {len(desc_boxes)} description boxes, expected 3")

        if desc_count == 3:
            total_score += 0.10
            print(f"PASS: Component 5 -- All 3 descriptions correct (0.10 pts)")
        elif desc_count > 0:
            partial = round(0.10 * desc_count / 3, 2)
            total_score += partial
            print(f"PARTIAL: Component 5 -- {desc_count}/3 correct ({partial:.2f} pts)")
        else:
            print(f"FAIL: Component 5 -- No description textboxes found with correct formatting")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
