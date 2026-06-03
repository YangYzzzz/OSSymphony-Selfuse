"""
FINAL REWARD SCRIPT - SUCCESS
Task: I'm touching up a huge deck in LibreOffice Impress, and the heading on slide 151 still fades into the background. How do I switch that title’s font color to Dark Blue 2 (#002F6C) and add a single underline so it finally jumps off the slide?
Generated: 2025-09-10 17:05:26
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
from pptx import Presentation
from pptx.dml.color import RGBColor

def _rgb_from_color_format(color):
    """Return an (r,g,b) tuple from a pptx.dml.color._ColorFormat object.
    Handles direct RGB as well as theme colors that have been modified.
    Returns None if RGB cannot be determined.
    """
    if color is None:
        return None

    # Direct RGB assignment
    if color.rgb is not None:
        return tuple(color.rgb)

    # Theme color case – attempt to resolve to RGB via theme part (best effort)
    try:
        theme_color = color.theme_color
        if theme_color is None:
            return None
        # theme_color.rgb returns RGBColor in recent python-pptx versions
        if hasattr(theme_color, "rgb") and theme_color.rgb is not None:
            return tuple(theme_color.rgb)
    except Exception:
        pass  # Fall through if theme resolution fails

    return None  # RGB could not be determined

def _runs_from_shape(shape):
    """Yield all text runs from a shape (including placeholders)."""
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            yield run

def verify_impress_heading_task(file_path):
    """Verify that on slide 151 the title font is Dark Blue 2 (#002F6C)
    and the text is single-underlined.
    Returns a progressive score between 0.0 and 1.0.
    """

    TARGET_RGB = (0x00, 0x2F, 0x6C)  # Dark Blue 2 (#002F6C)
    score = 0.0
    max_score = 1.0

    # --- 1. Load the presentation ------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File not found:", file_path)
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded presentation (slides: {len(prs.slides)})")
    except Exception as exc:
        print("✗ Error loading presentation:", exc)
        print("REWARD: 0.0")
        return 0.0

    # --- 2. Ensure slide 151 exists ---------------------------------------------
    if len(prs.slides) < 151:
        print(f"✗ Presentation has only {len(prs.slides)} slides – need at least 151")
        print("REWARD: 0.0")
        return 0.0

    slide_151 = prs.slides[150]  # zero-indexed

    # --- 3. Locate the title shape ----------------------------------------------
    title_shape = None

    # Prefer placeholder with type == TITLE (value 1)
    for shape in slide_151.shapes:
        if shape.has_text_frame and getattr(shape, "is_placeholder", False):
            if shape.placeholder_format.type == 1:  # TITLE placeholder
                title_shape = shape
                break

    # Fallback: first text shape on the slide
    if title_shape is None:
        for shape in slide_151.shapes:
            if shape.has_text_frame:
                title_shape = shape
                break

    if title_shape is None:
        print("✗ No text shape found on slide 151 – cannot verify title")
        print("REWARD: 0.0")
        return 0.0

    runs = list(_runs_from_shape(title_shape))
    if not runs:
        print("✗ Title shape contains no text runs")
        print("REWARD: 0.0")
        return 0.0

    # --- 4. Verify color ---------------------------------------------------------
    def run_has_target_color(run):
        rgb = _rgb_from_color_format(run.font.color)
        return rgb == TARGET_RGB

    all_color_ok = all(run_has_target_color(r) for r in runs)
    any_color_ok = any(run_has_target_color(r) for r in runs)

    if all_color_ok:
        print("✓ All title text is Dark Blue 2 (#002F6C)")
        score += 0.5
    elif any_color_ok:
        print("• Some, but not all, title text is Dark Blue 2 (#002F6C)")
        score += 0.25
    else:
        print("✗ Title text color is incorrect")

    # --- 5. Verify underline -----------------------------------------------------
    def run_is_underlined(run):
        ul = run.font.underline
        # pptx may store underline as bool or enumeration; any truthy / non-None value counts
        return bool(ul)

    all_ul_ok = all(run_is_underlined(r) for r in runs)
    any_ul_ok = any(run_is_underlined(r) for r in runs)

    if all_ul_ok:
        print("✓ All title text is underlined")
        score += 0.5
    elif any_ul_ok:
        print("• Some, but not all, title text is underlined")
        score += 0.25
    else:
        print("✗ Title text is not underlined")

    # --- 6. Final score ----------------------------------------------------------
    final_score = min(score, max_score)
    print(f"REWARD: {final_score}")
    return final_score

# -----------------------------------------------------------------------------
# Path to the presentation provided in the task context. Adjust if necessary.
FILE_PATH = "/home/user/im_touching_up_a_huge_deck_in_libreoffice_impress_and_the_heading_on_slide_151_still_fades_into_the__golden.pptx"

if __name__ == "__main__":
    verify_impress_heading_task(FILE_PATH)

