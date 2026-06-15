"""
Reward Script: Resize and reposition image on slide 3
Task ID: impress_teach_074
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30) - Image width is 4 inches
  Component 2 (0.30) - Image height is 3 inches
  Component 3 (0.20) - Image horizontally centered (left = 3 inches)
  Component 4 (0.20) - Image top edge at y = 4 inches
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_074'

# Tolerance: 0.5% relative tolerance for position/size checks
TOLERANCE = 0.005

def is_approx_equal(actual_emu, expected_inches):
    """Check if EMU value is approximately equal to expected inches value."""
    expected_emu = Inches(expected_inches)
    if expected_emu == 0:
        return actual_emu == 0
    return abs(actual_emu - expected_emu) / max(abs(actual_emu), abs(expected_emu)) <= TOLERANCE


def persist_app_state(domain):
    """Try to save any unsaved GUI state."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def find_image_on_slide(slide):
    """Find the first image shape on the slide."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return shape
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]
    image = find_image_on_slide(slide3)

    if image is None:
        print("FAIL: No image found on slide 3")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Image found - left={image.left/914400:.4f}in, top={image.top/914400:.4f}in, "
          f"width={image.width/914400:.4f}in, height={image.height/914400:.4f}in")

    # Component 1: Image width is 4 inches (0.30 points)
    try:
        if is_approx_equal(image.width, 4.0):
            print(f"PASS: Component 1 - Image width is ~4.0 inches ({image.width/914400:.4f} in) (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 1 - Expected width ~4.0 inches, found {image.width/914400:.4f} inches")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Image height is 3 inches (0.30 points)
    try:
        if is_approx_equal(image.height, 3.0):
            print(f"PASS: Component 2 - Image height is ~3.0 inches ({image.height/914400:.4f} in) (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 2 - Expected height ~3.0 inches, found {image.height/914400:.4f} inches")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Image horizontally centered (left = 3.0 inches for 10in slide, 4in image)
    # Center calculation: (slide_width - image_width) / 2
    try:
        slide_width = prs.slide_width
        expected_left_emu = (slide_width - image.width) // 2
        # Check that left position is approximately (slide_width - width) / 2
        # We use a slightly larger tolerance here since centering can have rounding
        if image.width > 0 and abs(image.left - expected_left_emu) <= Inches(0.1):
            print(f"PASS: Component 3 - Image horizontally centered (left={image.left/914400:.4f} in, "
                  f"expected ~{expected_left_emu/914400:.4f} in) (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 3 - Expected left ~{expected_left_emu/914400:.4f} in for centering, "
                  f"found {image.left/914400:.4f} in")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Image top edge at y = 4 inches (0.20 points)
    try:
        if is_approx_equal(image.top, 4.0):
            print(f"PASS: Component 4 - Image top at ~4.0 inches ({image.top/914400:.4f} in) (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 4 - Expected top ~4.0 inches, found {image.top/914400:.4f} inches")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
