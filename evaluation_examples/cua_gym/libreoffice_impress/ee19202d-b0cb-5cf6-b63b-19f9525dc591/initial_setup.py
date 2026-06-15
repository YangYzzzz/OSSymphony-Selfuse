"""
Initial Setup: Resize and reposition image on slide 3
Task ID: impress_teach_074
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_074'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_placeholder_image(path, width=400, height=300, label="Diagram"):
    """Create a simple placeholder image using PIL."""
    try:
        from PIL import Image, ImageDraw, ImageFont
        img = Image.new('RGB', (width, height), color=(220, 235, 250))
        draw = ImageDraw.Draw(img)
        # Draw border
        draw.rectangle([0, 0, width-1, height-1], outline=(70, 130, 180), width=3)
        # Draw diagonal lines for visual interest
        draw.line([(0, 0), (width, height)], fill=(180, 200, 230), width=1)
        draw.line([(width, 0), (0, height)], fill=(180, 200, 230), width=1)
        # Add label text
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
        except (IOError, OSError):
            font = ImageFont.load_default()
        bbox = draw.textbbox((0, 0), label, font=font)
        tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
        draw.text(((width - tw) / 2, (height - th) / 2), label, fill=(50, 80, 120), font=font)
        img.save(path)
    except ImportError:
        # Fallback: create a minimal 1x1 PNG
        import struct, zlib
        def create_minimal_png(path):
            sig = b'\x89PNG\r\n\x1a\n'
            ihdr_data = struct.pack('>IIBBBBB', width, height, 8, 2, 0, 0, 0)
            ihdr_crc = zlib.crc32(b'IHDR' + ihdr_data) & 0xffffffff
            ihdr = struct.pack('>I', 13) + b'IHDR' + ihdr_data + struct.pack('>I', ihdr_crc)
            raw = b''
            for _ in range(height):
                raw += b'\x00' + b'\xdc\xeb\xfa' * width
            compressed = zlib.compress(raw)
            idat_crc = zlib.crc32(b'IDAT' + compressed) & 0xffffffff
            idat = struct.pack('>I', len(compressed)) + b'IDAT' + compressed + struct.pack('>I', idat_crc)
            iend_crc = zlib.crc32(b'IEND') & 0xffffffff
            iend = struct.pack('>I', 0) + b'IEND' + struct.pack('>I', iend_crc)
            with open(path, 'wb') as f:
                f.write(sig + ihdr + idat + iend)
        create_minimal_png(path)


def create_initial():
    prs = Presentation()
    # Standard 10 x 7.5 inch slide
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Introduction to Data Visualization"
    slide1.placeholders[1].text = "Dr. Elena Rodriguez\nComputer Science Department\nSpring 2026"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Lecture Outline"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Why Visualization Matters"
    items = [
        "Types of Charts and Their Use Cases",
        "Choosing the Right Visualization",
        "Color Theory in Data Presentation",
        "Hands-On Exercise: Building Dashboards",
        "Summary and Next Steps",
    ]
    for item in items:
        p = body2.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 3: Key slide with image (top-left, 2x1.5 inches) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title textbox
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Chart Type Comparison"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x6E)

    # Create a placeholder image and add it at top-left corner, 2x1.5 inches
    img_path = os.path.join(tempfile.gettempdir(), 'chart_comparison.png')
    create_placeholder_image(img_path, 600, 450, "Chart Types")
    slide3.shapes.add_picture(img_path, Inches(0.2), Inches(0.2), Inches(2), Inches(1.5))

    # Add explanatory text below the image area
    txBox2 = slide3.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(3))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Bar charts are ideal for comparing quantities across categories."
    run2 = p2.runs[0]
    run2.font.size = Pt(16)
    p3 = tf2.add_paragraph()
    p3.text = "Line charts work best for showing trends over continuous intervals."
    for r in p3.runs:
        r.font.size = Pt(16)
    p4 = tf2.add_paragraph()
    p4.text = "Pie charts should only be used for parts-of-a-whole with fewer than 6 slices."
    for r in p4.runs:
        r.font.size = Pt(16)

    # --- Slide 4: Color Theory ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Color Theory for Data"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Sequential: Light to dark for ordered data"
    for line in [
        "Diverging: Two hues for positive/negative deviation",
        "Categorical: Distinct hues for unrelated groups",
        "Avoid red-green combinations (colorblind users)",
        "Keep palettes to 5-7 distinct colors maximum",
    ]:
        p = body4.add_paragraph()
        p.text = line

    # --- Slide 5: Best Practices ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Visualization Best Practices"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Label axes clearly with units"
    for line in [
        "Use consistent scales across comparisons",
        "Minimize chart junk and decorative elements",
        "Highlight the key insight, not all the data",
        "Test readability at presentation resolution",
    ]:
        p = body5.add_paragraph()
        p.text = line

    # --- Slide 6: Summary ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Key Takeaways"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Match chart type to the data relationship"
    for line in [
        "Use color deliberately, not decoratively",
        "Always consider your audience and context",
        "Practice with the dashboard exercise this week",
        "Next lecture: Interactive Visualization with Plotly",
    ]:
        p = body6.add_paragraph()
        p.text = line

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Clean up temp image
    try:
        os.remove(img_path)
    except OSError:
        pass

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
