"""
Reward Script: Apply underline and #800000 (dark red) to title, body bullets, and table cells on slide 4
Task ID: osworld_impress_underline_darkred_table_007
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Title textbox (TextBox 2, 'Product Line Analysis') has underline=True and color=800000
  Component 2 (0.40): All 3 body bullets in TextBox 3 have underline=True and color=800000
  Component 3 (0.30): All 16 cells in the 4x4 table have underline=True and color=800000 in all runs
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_underline_darkred_table_007'

DARK_RED = '800000'  # Expected hex color — pptx returns uppercase RGB string


def normalize_color(rgb_str):
    """Normalize hex color string to uppercase for comparison."""
    return str(rgb_str).upper() if rgb_str else None


def run_passes_formatting(run):
    """
    Returns (underline_ok, color_ok) for a given run.
    Checks: underline is True AND color is 800000.
    """
    underline_ok = run.font.underline is True
    color_ok = False
    try:
        if run.font.color.type is not None:
            color_ok = normalize_color(run.font.color.rgb) == DARK_RED.upper()
    except Exception:
        color_ok = False
    return underline_ok, color_ok


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except Exception as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition gate: confirm there are at least 4 slides
    if len(prs.slides) < 4:
        print(f"CRITICAL: Expected at least 4 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)

    # Identify the relevant shapes on slide 4:
    #   Shape 0 (Title 1): empty placeholder — not scored
    #   Shape 1 (TextBox 2): title text "Product Line Analysis"
    #   Shape 2 (TextBox 3): body bullets (3 items)
    #   Shape 3 (Table 4): 4x4 data table

    title_shape = None
    body_shape = None
    table_shape = None

    for shape in slide.shapes:
        if shape.name == 'TextBox 2' and shape.has_text_frame:
            title_shape = shape
        elif shape.name == 'TextBox 3' and shape.has_text_frame:
            body_shape = shape
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_shape = shape

    # --- Component 1: Title textbox has underline=True and color=800000 (0.30 points) ---
    # Verifies "Product Line Analysis" run is formatted with underline + dark red.
    # In the initial state, this run has underline=False and color=000000.
    try:
        if title_shape is None:
            print("FAIL: Component 1 — TextBox 2 (title) not found on slide 4")
        else:
            runs = [r for para in title_shape.text_frame.paragraphs
                    for r in para.runs if (r.text or "").strip()]
            if not runs:
                print("FAIL: Component 1 — No non-empty runs in title TextBox 2")
            else:
                fail_count = 0
                fail_details = []
                for run in runs:
                    underline_ok, color_ok = run_passes_formatting(run)
                    if not underline_ok:
                        fail_count += 1
                        fail_details.append(f"'{run.text[:20]}': underline={run.font.underline}")
                    if not color_ok:
                        fail_count += 1
                        actual_c = 'none'
                        try:
                            if run.font.color.type is not None:
                                actual_c = str(run.font.color.rgb)
                        except Exception:
                            pass
                        fail_details.append(f"'{run.text[:20]}': color={actual_c}")

                if fail_count == 0:
                    print(f"PASS: Component 1 — Title 'Product Line Analysis' has underline=True and color=800000 (0.30 pts)")
                    total_score += 0.30
                else:
                    print(f"FAIL: Component 1 — Title formatting issues ({fail_count}): {fail_details}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # --- Component 2: All body bullets have underline=True and color=800000 (0.40 points) ---
    # Verifies all 3 bullet paragraphs in TextBox 3 are formatted correctly.
    # In the initial state all runs have underline=False and color=000000.
    try:
        if body_shape is None:
            print("FAIL: Component 2 — TextBox 3 (body bullets) not found on slide 4")
        else:
            fail_count = 0
            total_runs_checked = 0
            fail_details = []

            for para in body_shape.text_frame.paragraphs:
                nonempty_runs = [r for r in para.runs if (r.text or "").strip()]
                for run in nonempty_runs:
                    total_runs_checked += 1
                    underline_ok, color_ok = run_passes_formatting(run)
                    if not underline_ok:
                        fail_count += 1
                        fail_details.append(f"'{run.text[:30]}': underline={run.font.underline}")
                    if not color_ok:
                        fail_count += 1
                        actual_c = 'none'
                        try:
                            if run.font.color.type is not None:
                                actual_c = str(run.font.color.rgb)
                        except Exception:
                            pass
                        fail_details.append(f"'{run.text[:30]}': color={actual_c}")

            if total_runs_checked == 0:
                print("FAIL: Component 2 — No non-empty runs in body TextBox 3")
            elif fail_count == 0:
                print(f"PASS: Component 2 — All {total_runs_checked} body bullet runs have underline=True and color=800000 (0.40 pts)")
                total_score += 0.40
            else:
                print(f"FAIL: Component 2 — Body bullet formatting issues ({fail_count} problems): {fail_details[:5]}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # --- Component 3: All 16 table cells have underline=True and color=800000 (0.30 points) ---
    # Verifies every cell in the 4x4 table is formatted with underline + dark red.
    # In the initial state all table cell runs have underline=False and color=000000.
    try:
        if table_shape is None:
            print("FAIL: Component 3 — Table shape not found on slide 4")
        else:
            table = table_shape.table
            rows = len(table.rows)
            cols = len(table.columns)
            fail_count = 0
            total_cell_runs = 0
            fail_details = []

            for r in range(rows):
                for c in range(cols):
                    cell = table.cell(r, c)
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            if not (run.text or "").strip():
                                continue
                            total_cell_runs += 1
                            underline_ok, color_ok = run_passes_formatting(run)
                            if not underline_ok:
                                fail_count += 1
                                fail_details.append(f"cell[{r},{c}] '{run.text[:12]}': underline={run.font.underline}")
                            if not color_ok:
                                fail_count += 1
                                actual_c = 'none'
                                try:
                                    if run.font.color.type is not None:
                                        actual_c = str(run.font.color.rgb)
                                except Exception:
                                    pass
                                fail_details.append(f"cell[{r},{c}] '{run.text[:12]}': color={actual_c}")

            if total_cell_runs == 0:
                print("FAIL: Component 3 — No non-empty runs found in table cells")
            elif fail_count == 0:
                print(f"PASS: Component 3 — All {rows}x{cols} table cells ({total_cell_runs} runs) have underline=True and color=800000 (0.30 pts)")
                total_score += 0.30
            else:
                print(f"FAIL: Component 3 — Table cell formatting issues ({fail_count} problems): {fail_details[:8]}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the VM environment
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
