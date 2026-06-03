"""
Initial Setup: 8-slide data summary deck with slide 4 containing title, body bullets, and a 4x4 table in black text.
Task ID: osworld_impress_underline_darkred_table_007
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_underline_darkred_table_007'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

BLACK = RGBColor(0x00, 0x00, 0x00)


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_run_style(run, bold=False, underline=False, color=BLACK, size_pt=18):
    run.font.bold = bold
    run.font.underline = underline
    run.font.color.rgb = color
    run.font.size = Pt(size_pt)


def add_text_to_tf(tf, paragraphs_data, default_size=18):
    """paragraphs_data: list of (text, bold, underline, color, size)"""
    for i, (text, bold, underline, color, size) in enumerate(paragraphs_data):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.text = text
        for run in para.runs:
            set_run_style(run, bold=bold, underline=underline, color=color, size_pt=size)


def create_initial():
    prs = Presentation()
    # Standard 10x7.5 inch slide
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_layouts = prs.slide_layouts
    # Use Title+Content layout for most slides
    title_content_layout = slide_layouts[1]
    title_only_layout = slide_layouts[5]  # blank

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(slide_layouts[0])
    slide1.shapes.title.text = "Q1 2025 Performance Summary"
    slide1.placeholders[1].text = "Annual Business Review\nFiscal Year Data"

    # ---- Slide 2: Overview ----
    slide2 = prs.slides.add_slide(title_content_layout)
    slide2.shapes.title.text = "Executive Overview"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Revenue increased by 12% year-over-year"
    tf2.add_paragraph().text = "Operating costs reduced by 7% through efficiency initiatives"
    tf2.add_paragraph().text = "Three new product lines launched successfully"
    tf2.add_paragraph().text = "Customer satisfaction score: 4.7 / 5.0"
    for para in tf2.paragraphs:
        for run in para.runs:
            run.font.color.rgb = BLACK
            run.font.size = Pt(18)

    # ---- Slide 3: Regional Performance ----
    slide3 = prs.slides.add_slide(title_content_layout)
    slide3.shapes.title.text = "Regional Performance"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "North America: $42.3M (+15%)"
    tf3.add_paragraph().text = "Europe: $27.1M (+8%)"
    tf3.add_paragraph().text = "Asia Pacific: $18.6M (+22%)"
    tf3.add_paragraph().text = "Latin America: $9.4M (+5%)"
    for para in tf3.paragraphs:
        for run in para.runs:
            run.font.color.rgb = BLACK
            run.font.size = Pt(18)

    # ---- Slide 4: Product Line Analysis (KEY SLIDE) ----
    # Title + body textbox with 3 bullets + 4x4 table, all in black text, NO underline
    slide4 = prs.slides.add_slide(title_only_layout)

    # Title
    title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.9))
    tf_title = title_box.text_frame
    tf_title.text = "Product Line Analysis"
    run_t = tf_title.paragraphs[0].runs[0]
    run_t.font.bold = True
    run_t.font.underline = False
    run_t.font.color.rgb = BLACK
    run_t.font.size = Pt(28)

    # Body textbox with 3 bullet items
    body_box = slide4.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(2.0))
    tf_body = body_box.text_frame
    tf_body.word_wrap = True
    tf_body.text = "Alpha Series contributed 38% of total product revenue"
    p2 = tf_body.add_paragraph()
    p2.text = "Beta Platform unit sales grew by 29% compared to prior quarter"
    p3 = tf_body.add_paragraph()
    p3.text = "Gamma Tools division exceeded annual targets by $1.2M"
    for para in tf_body.paragraphs:
        for run in para.runs:
            run.font.underline = False
            run.font.color.rgb = BLACK
            run.font.size = Pt(18)

    # 4x4 table (header row + 3 data rows, 4 columns)
    rows, cols = 4, 4
    table_shape = slide4.shapes.add_table(
        rows, cols,
        Inches(0.5), Inches(3.6),
        Inches(9.0), Inches(2.8)
    )
    table = table_shape.table

    cell_data = [
        ["Product",    "Q1 Sales ($)",  "Q2 Sales ($)",  "Growth (%)"],
        ["Alpha",      "1,245,000",     "1,720,400",     "+38.2"],
        ["Beta",       "987,300",       "1,272,800",     "+28.9"],
        ["Gamma",      "654,100",       "731,500",       "+11.8"],
    ]

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = cell_data[r][c]
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.bold = (r == 0)
                    run.font.underline = False
                    run.font.color.rgb = BLACK
                    run.font.size = Pt(16)

    # ---- Slide 5: Cost Structure ----
    slide5 = prs.slides.add_slide(title_content_layout)
    slide5.shapes.title.text = "Cost Structure Breakdown"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Personnel costs: 45% of total operating budget"
    tf5.add_paragraph().text = "Infrastructure and IT: 18% allocation"
    tf5.add_paragraph().text = "Marketing and sales: 22% year-to-date"
    tf5.add_paragraph().text = "R&D investment: 15% in line with strategy"
    for para in tf5.paragraphs:
        for run in para.runs:
            run.font.color.rgb = BLACK
            run.font.size = Pt(18)

    # ---- Slide 6: Customer Insights ----
    slide6 = prs.slides.add_slide(title_content_layout)
    slide6.shapes.title.text = "Customer Insights"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Net Promoter Score improved to 72 (up from 65)"
    tf6.add_paragraph().text = "Customer retention rate: 91.3%"
    tf6.add_paragraph().text = "Support ticket resolution time reduced by 34%"
    tf6.add_paragraph().text = "New customer acquisitions: 1,840 accounts"
    for para in tf6.paragraphs:
        for run in para.runs:
            run.font.color.rgb = BLACK
            run.font.size = Pt(18)

    # ---- Slide 7: Strategic Initiatives ----
    slide7 = prs.slides.add_slide(title_content_layout)
    slide7.shapes.title.text = "Strategic Initiatives Q2 2025"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "Launch Delta product line in APAC markets by April"
    tf7.add_paragraph().text = "Expand enterprise sales team by 12 headcount"
    tf7.add_paragraph().text = "Complete ERP migration to cloud platform"
    tf7.add_paragraph().text = "Achieve ISO 27001 certification by June"
    for para in tf7.paragraphs:
        for run in para.runs:
            run.font.color.rgb = BLACK
            run.font.size = Pt(18)

    # ---- Slide 8: Closing ----
    slide8 = prs.slides.add_slide(slide_layouts[0])
    slide8.shapes.title.text = "Thank You"
    slide8.placeholders[1].text = "Questions & Discussion\ncontact@company.com"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
