"""
FINAL REWARD SCRIPT - SUCCESS
Task: I need to make a particular shape stand out on my slide by giving it a special look. How can I set up a drawing style called 'Highlight-Box' with a yellow fill (#FFFF00) and a red outline (#FF0000) in Impress? I'm not quite sure how to apply that style to the shapes I've selected.
Generated: 2025-08-07 11:35:05
Status: success
Model: o4-mini
Total Steps: 2
"""

import os
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.dml.color import RGBColor

def verify_highlight_box(file_path):
    print("Starting verification of Highlight-Box style application...")
    total_score = 0.0
    max_score = 1.0

    # Requirement 1: File exists
    try:
        if os.path.exists(file_path):
            print(f"✓ File exists: {file_path} (0.2 points)")
            total_score += 0.2
        else:
            print(f"✗ File not found: {file_path}")
            print(f"REWARD: {total_score}")
            return total_score
    except Exception as e:
        print(f"✗ Error checking file existence: {e}")
        print(f"REWARD: {total_score}")
        return total_score

    # Requirement 2: Load presentation and check slide count
    try:
        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        if slide_count > 0:
            print(f"✓ Presentation loaded with {slide_count} slides (0.3 points)")
            total_score += 0.3
        else:
            print("✗ Presentation has no slides (0 points)")
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        print(f"REWARD: {total_score}")
        return total_score

    # Requirement 3: Check for shape with yellow fill and red outline
    found = False
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            try:
                fill = shape.fill
                # Check for solid yellow fill (#FFFF00)
                if fill and fill.type == MSO_FILL.SOLID and fill.fore_color.rgb == RGBColor(255, 255, 0):
                    line_fill = shape.line.fill
                    # Check for solid red outline (#FF0000)
                    if line_fill and line_fill.type == MSO_FILL.SOLID and line_fill.fore_color.rgb == RGBColor(255, 0, 0):
                        print(f"✓ Found shape with yellow fill and red outline on slide {slide_idx + 1}")
                        found = True
                        break
            except Exception:
                # Some shapes lack fill/line attributes; skip
                continue
        if found:
            break

    if found:
        print("✓ Highlight-Box style applied to a shape (0.5 points)")
        total_score += 0.5
    else:
        print("✗ No shape found with yellow fill and red outline (0 points)")

    final_score = min(total_score, max_score)
    print(f"REWARD: {final_score}")
    return final_score

if __name__ == '__main__':
    # Path to the presentation file to verify
    file_path = '/home/user/i_need_to_make_a_particular_shape_stand_out_on_my_slide_by_giving_it_a_special_look_how_can_i_set_up.pptx'
    verify_highlight_box(file_path)
