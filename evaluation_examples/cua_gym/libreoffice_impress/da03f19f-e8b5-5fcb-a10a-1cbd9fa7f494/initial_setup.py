"""
Initial Setup: Presentation with animations on slide 3.
Task ID: impress_ma_068
Domain: libreoffice_impress
Slide 3 animations in order: (1) Title - Fade In, (2) Bullet List - Appear, (3) Chart - Wipe
"""

import os
import shlex
import subprocess
import time
import copy
import zipfile
import shutil
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_068'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # === Slide 1: Title Slide ===
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q3 2025 Financial Review"
    slide1.placeholders[1].text = "Prepared by the Strategy & Finance Team\nSeptember 2025"

    # === Slide 2: Overview ===
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Revenue performance and trends"
    items = [
        "Cost analysis by department",
        "Regional growth breakdown",
        "Strategic initiatives update",
        "Q4 outlook and projections",
    ]
    for item in items:
        p = tf2.add_paragraph()
        p.text = item
        p.level = 0

    # === Slide 3: The target slide with animations ===
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Shape 1: Title text box
    title_box = slide3.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(8.4), Inches(1.0))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "Regional Revenue Comparison"
    p_title.alignment = PP_ALIGN.CENTER
    run_title = p_title.runs[0]
    run_title.font.size = Pt(28)
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Shape 2: Bullet list text box
    bullet_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(4.0), Inches(4.5))
    tf_bullet = bullet_box.text_frame
    tf_bullet.word_wrap = True
    bullet_items = [
        ("North America", "+12% YoY growth driven by enterprise sales"),
        ("Europe", "+8% with strong performance in DACH region"),
        ("Asia Pacific", "+22% led by expansion in Japan and India"),
        ("Latin America", "+5% despite currency headwinds"),
    ]
    for i, (region, detail) in enumerate(bullet_items):
        if i == 0:
            p = tf_bullet.paragraphs[0]
        else:
            p = tf_bullet.add_paragraph()
        p.text = region
        p.level = 0
        run = p.runs[0]
        run.font.bold = True
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x2E, 0x75, 0xB6)

        p_detail = tf_bullet.add_paragraph()
        p_detail.text = detail
        p_detail.level = 1
        run_d = p_detail.runs[0]
        run_d.font.size = Pt(12)
        run_d.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # Shape 3: Chart placeholder (a grouped shape simulating a bar chart)
    chart_box = slide3.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(4.5))
    tf_chart = chart_box.text_frame
    tf_chart.word_wrap = True
    p_chart = tf_chart.paragraphs[0]
    p_chart.text = "[Revenue Chart - Bar Graph]"
    p_chart.alignment = PP_ALIGN.CENTER
    run_chart = p_chart.runs[0]
    run_chart.font.size = Pt(14)
    run_chart.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    # Add some bar-like data text
    chart_data = [
        "NA: $45.2M  ████████████",
        "EU: $31.8M  ████████",
        "APAC: $28.5M  ███████",
        "LATAM: $12.1M  ███",
    ]
    for line in chart_data:
        p_c = tf_chart.add_paragraph()
        p_c.text = line
        run_c = p_c.runs[0]
        run_c.font.size = Pt(11)
        run_c.font.name = "Courier New"
        run_c.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # === Slide 4: Details ===
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Cost Optimization Results"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "IT infrastructure: 18% reduction through cloud migration"
    for item in [
        "Marketing spend: Reallocated 25% to digital channels",
        "Operations: Automated 40% of manual processes",
        "Travel: 30% savings via hybrid meeting policy",
    ]:
        p = tf4.add_paragraph()
        p.text = item

    # === Slide 5: Summary ===
    slide5 = prs.slides.add_slide(prs.slide_layouts[0])
    slide5.shapes.title.text = "Key Takeaways & Next Steps"
    slide5.placeholders[1].text = "Strong momentum across all regions\nFocus on APAC expansion in Q4"

    # Save the presentation first (without animations)
    prs.save(OUTPUT)

    # Now add animations via XML manipulation
    add_animations_to_slide3(OUTPUT, title_box, bullet_box, chart_box, slide3)

    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


def add_animations_to_slide3(pptx_path, title_shape, bullet_shape, chart_shape, slide):
    """Add animations to slide 3 via direct XML manipulation of the saved file.

    Initial order:
      1. Title - Fade In (On Click)
      2. Bullet List - Appear (On Click)
      3. Chart - Wipe (On Click)
    """
    # Get shape IDs from the python-pptx objects
    title_sp_id = title_shape.shape_id
    bullet_sp_id = bullet_shape.shape_id
    chart_sp_id = chart_shape.shape_id

    # Build the timing XML for slide 3
    # Namespaces
    nsmap = {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    }

    def make_anim_seq(animations):
        """Create the full p:timing XML element with an animation sequence.

        animations: list of (shape_id, shape_name, effect_type, anim_idx)
        effect_type: 'fade', 'appear', 'wipe'
        """
        # We'll build the XML string and parse it - much cleaner for complex structures
        anim_entries = []
        for idx, (sp_id, sp_name, effect_type, _) in enumerate(animations):
            delay = "0" if idx == 0 else "0"

            if effect_type == 'fade':
                anim_effect = f'''
                    <p:animEffect transition="in" filter="fade">
                      <p:cBhvr>
                        <p:cTn id="{10 + idx*10 + 5}" dur="500" />
                        <p:tgtEl>
                          <p:spTgt spid="{sp_id}" />
                        </p:tgtEl>
                      </p:cBhvr>
                    </p:animEffect>'''
                set_elem = f'''
                    <p:set>
                      <p:cBhvr>
                        <p:cTn id="{10 + idx*10 + 6}" dur="1" fill="hold">
                          <p:stCondLst>
                            <p:cond delay="0" />
                          </p:stCondLst>
                        </p:cTn>
                        <p:tgtEl>
                          <p:spTgt spid="{sp_id}" />
                        </p:tgtEl>
                        <p:attrNameLst>
                          <p:attrName>style.visibility</p:attrName>
                        </p:attrNameLst>
                      </p:cBhvr>
                      <p:to>
                        <p:strVal val="visible" />
                      </p:to>
                    </p:set>'''
            elif effect_type == 'appear':
                anim_effect = ''
                set_elem = f'''
                    <p:set>
                      <p:cBhvr>
                        <p:cTn id="{10 + idx*10 + 5}" dur="1" fill="hold">
                          <p:stCondLst>
                            <p:cond delay="0" />
                          </p:stCondLst>
                        </p:cTn>
                        <p:tgtEl>
                          <p:spTgt spid="{sp_id}" />
                        </p:tgtEl>
                        <p:attrNameLst>
                          <p:attrName>style.visibility</p:attrName>
                        </p:attrNameLst>
                      </p:cBhvr>
                      <p:to>
                        <p:strVal val="visible" />
                      </p:to>
                    </p:set>'''
            elif effect_type == 'wipe':
                anim_effect = f'''
                    <p:animEffect transition="in" filter="wipe(down)">
                      <p:cBhvr>
                        <p:cTn id="{10 + idx*10 + 5}" dur="500" />
                        <p:tgtEl>
                          <p:spTgt spid="{sp_id}" />
                        </p:tgtEl>
                      </p:cBhvr>
                    </p:animEffect>'''
                set_elem = f'''
                    <p:set>
                      <p:cBhvr>
                        <p:cTn id="{10 + idx*10 + 6}" dur="1" fill="hold">
                          <p:stCondLst>
                            <p:cond delay="0" />
                          </p:stCondLst>
                        </p:cTn>
                        <p:tgtEl>
                          <p:spTgt spid="{sp_id}" />
                        </p:tgtEl>
                        <p:attrNameLst>
                          <p:attrName>style.visibility</p:attrName>
                        </p:attrNameLst>
                      </p:cBhvr>
                      <p:to>
                        <p:strVal val="visible" />
                      </p:to>
                    </p:set>'''

            # Each animation is a separate click sequence entry
            if idx == 0:
                start_cond = '<p:stCondLst><p:cond delay="0" /></p:stCondLst>'
            else:
                start_cond = '<p:stCondLst><p:cond delay="0" /></p:stCondLst>'

            child_tn_children = set_elem + anim_effect

            anim_entry = f'''
              <p:par>
                <p:cTn id="{10 + idx*10 + 1}" fill="hold">
                  <p:stCondLst>
                    <p:cond delay="0" />
                  </p:stCondLst>
                  <p:childTnLst>
                    <p:par>
                      <p:cTn id="{10 + idx*10 + 2}" fill="hold">
                        <p:stCondLst>
                          <p:cond delay="0" />
                        </p:stCondLst>
                        <p:childTnLst>
                          <p:par>
                            <p:cTn id="{10 + idx*10 + 3}" presetID="{_preset_id(effect_type)}" presetClass="entr" presetSubtype="{_preset_subtype(effect_type)}" fill="hold" nodeType="clickEffect">
                              <p:stCondLst>
                                <p:cond delay="0" />
                              </p:stCondLst>
                              <p:childTnLst>
                                {child_tn_children}
                              </p:childTnLst>
                            </p:cTn>
                          </p:par>
                        </p:childTnLst>
                      </p:cTn>
                    </p:par>
                  </p:childTnLst>
                </p:cTn>
              </p:par>'''
            anim_entries.append(anim_entry)

        # Build the complete timing XML
        seq_entries = "\n".join(anim_entries)

        timing_xml = f'''<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                                    xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <p:tnLst>
            <p:par>
              <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
                <p:childTnLst>
                  <p:seq concurrent="1" nextAc="seek">
                    <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                      <p:childTnLst>
                        {seq_entries}
                      </p:childTnLst>
                    </p:cTn>
                    <p:prevCondLst>
                      <p:cond evt="onPrev" delay="0">
                        <p:tgtEl>
                          <p:sldTgt />
                        </p:tgtEl>
                      </p:cond>
                    </p:prevCondLst>
                    <p:nextCondLst>
                      <p:cond evt="onNext" delay="0">
                        <p:tgtEl>
                          <p:sldTgt />
                        </p:tgtEl>
                      </p:cond>
                    </p:nextCondLst>
                  </p:seq>
                </p:childTnLst>
              </p:cTn>
            </p:par>
          </p:tnLst>
        </p:timing>'''

        return timing_xml

    # Initial animation order: Title (Fade), Bullet (Appear), Chart (Wipe)
    animations = [
        (title_sp_id, "Title", "fade", 0),
        (bullet_sp_id, "BulletList", "appear", 1),
        (chart_sp_id, "Chart", "wipe", 2),
    ]

    timing_xml = make_anim_seq(animations)

    # Now inject this into the saved pptx file (slide3.xml)
    import tempfile
    tmp_path = pptx_path + '.tmp'

    with zipfile.ZipFile(pptx_path, 'r') as zin:
        with zipfile.ZipFile(tmp_path, 'w') as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == 'ppt/slides/slide3.xml':
                    # Parse and modify
                    root = etree.fromstring(data)
                    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

                    # Remove existing timing if any
                    for old_timing in root.findall('.//p:timing', ns):
                        old_timing.getparent().remove(old_timing)

                    # Parse new timing element
                    timing_el = etree.fromstring(timing_xml.encode('utf-8'))

                    # Append to slide root
                    root.append(timing_el)

                    data = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)

                zout.writestr(item, data)

    shutil.move(tmp_path, pptx_path)
    print(f'Animations added to slide 3')


def _preset_id(effect_type):
    """Return the OOXML preset ID for common animation effects."""
    return {'fade': 10, 'appear': 1, 'wipe': 22}[effect_type]


def _preset_subtype(effect_type):
    """Return the preset subtype."""
    return {'fade': 0, 'appear': 0, 'wipe': 4}[effect_type]  # wipe down = subtype 4


create_initial()
