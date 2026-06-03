"""
Reward Script: Split-content debate slide with FOR/AGAINST headers, bullets, dividing line, and VS circle
Task ID: impress_stu_094
Domain: libreoffice_impress
Scoring:
  Component 1 (0.15): Vertical dividing line exists on slide 5
  Component 2 (0.20): FOR header with green color (#27AE60) and bold
  Component 3 (0.15): 4 supporting argument bullets on left side (~14pt)
  Component 4 (0.20): AGAINST header with red color (#C0392B) and bold
  Component 5 (0.15): 4 counter-argument bullets on right side (~14pt)
  Component 6 (0.15): VS circle/oval shape with gold (#F1C40F) fill and bold text
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_094'


def get_shape_fill_rgb(shape):
    """Get the solid fill RGB color of a shape, or None."""
    try:
        fill = shape.fill
        if fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def get_run_color(run):
    """Get the RGB color string of a run's font, or None."""
    try:
        if run.font.color.type is not None:
            return str(run.font.color.rgb)
    except Exception:
        pass
    return None


def get_all_text_from_shape(shape):
    """Get concatenated text from a shape's text frame."""
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text.strip()


def count_nonempty_paragraphs(shape):
    """Count paragraphs with non-empty text in a shape."""
    if not shape.has_text_frame:
        return 0
    count = 0
    for para in shape.text_frame.paragraphs:
        if para.text.strip():
            count += 1
    return count


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # Slide 5 (0-indexed)
    shapes = list(slide.shapes)
    slide_width = prs.slide_width  # ~9144000 EMU for standard 10"

    # Classify shapes on slide 5 by type and position
    line_shapes = []
    text_shapes = []
    auto_shapes = []

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            line_shapes.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            auto_shapes.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            text_shapes.append(shape)
        # Skip PLACEHOLDERs (they are pre-existing)

    # Component 1: Vertical dividing line exists (0.15 points)
    # A vertical line should have width ~0 (or very small) and significant height, positioned near center
    try:
        found_vertical_line = False
        center_x = slide_width // 2  # ~4572000
        tolerance = slide_width * 0.15  # 15% tolerance from center

        for shape in line_shapes:
            # A vertical line: width is 0 or very small, height is significant
            is_vertical = (shape.width <= Emu(50000)) and (shape.height > Emu(500000))
            # Check near center
            near_center = abs(shape.left - center_x) < tolerance
            if is_vertical and near_center:
                found_vertical_line = True
                print(f"PASS: Component 1 — Vertical line found at x={shape.left}, height={shape.height} (0.15 pts)")
                total_score += 0.15
                break

        if not found_vertical_line:
            # Also check freeform/connector shapes that may serve as lines
            for shape in shapes:
                if shape.shape_type in (MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM):
                    is_vertical = (shape.width <= Emu(50000)) and (shape.height > Emu(500000))
                    near_center = abs(shape.left - center_x) < tolerance
                    if is_vertical and near_center:
                        found_vertical_line = True
                        print(f"PASS: Component 1 — Vertical line/connector found at x={shape.left} (0.15 pts)")
                        total_score += 0.15
                        break

        if not found_vertical_line:
            print(f"FAIL: Component 1 — No vertical dividing line found near center of slide 5")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: FOR header with green color (#27AE60) and bold (0.20 points)
    try:
        found_for_header = False
        for_header_score = 0.0

        for shape in text_shapes:
            text = get_all_text_from_shape(shape)
            if 'FOR' in text.upper() and 'UNIVERSAL BASIC INCOME' in text.upper() and 'AGAINST' not in text.upper():
                # Found the FOR header text box
                # Check it's on the left side (left position < center)
                if shape.left < center_x:
                    found_for_header = True
                    # Check color and bold
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                color = get_run_color(run)
                                is_bold = run.font.bold is True
                                is_green = color is not None and color.upper() == '27AE60'
                                # Check size ~22pt (279400 EMU = 22pt)
                                size_ok = run.font.size is not None and abs(run.font.size - Pt(22)) < Pt(3)

                                if is_green and is_bold:
                                    for_header_score = 0.20
                                    print(f"PASS: Component 2 — FOR header: green={is_green}, bold={is_bold}, size={run.font.size} (0.20 pts)")
                                elif is_green or is_bold:
                                    for_header_score = 0.10
                                    print(f"PARTIAL: Component 2 — FOR header: green={is_green}, bold={is_bold} (0.10 pts)")
                                else:
                                    print(f"FAIL: Component 2 — FOR header found but color={color}, bold={run.font.bold}")
                                break
                        break
                    break

        if not found_for_header:
            print(f"FAIL: Component 2 — No 'FOR Universal Basic Income' header found on left side of slide 5")

        total_score += for_header_score
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: 4 supporting argument bullets on left side (~14pt) (0.15 points)
    try:
        found_for_bullets = False
        for_bullet_score = 0.0

        for shape in text_shapes:
            text = get_all_text_from_shape(shape)
            # The bullet text box should be on the left side and NOT contain the header text
            if shape.left < center_x and 'FOR' not in text.upper().split('\n')[0][:10]:
                # Check if it's a multi-paragraph text box (bullets)
                para_count = count_nonempty_paragraphs(shape)
                if para_count >= 3:  # Looking for ~4 bullets
                    found_for_bullets = True
                    # Check font size ~14pt (177800 EMU = 14pt)
                    size_ok = False
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size is not None and abs(run.font.size - Pt(14)) < Pt(3):
                                size_ok = True
                                break
                        if size_ok:
                            break

                    if para_count >= 4 and size_ok:
                        for_bullet_score = 0.15
                        print(f"PASS: Component 3 — {para_count} FOR bullets found at ~14pt (0.15 pts)")
                    elif para_count >= 4:
                        for_bullet_score = 0.10
                        print(f"PARTIAL: Component 3 — {para_count} FOR bullets found but size not 14pt (0.10 pts)")
                    elif para_count >= 2:
                        for_bullet_score = 0.07
                        print(f"PARTIAL: Component 3 — Only {para_count} FOR bullets found (0.07 pts)")
                    else:
                        print(f"FAIL: Component 3 — Only {para_count} FOR bullets")
                    break

        if not found_for_bullets:
            print(f"FAIL: Component 3 — No supporting argument bullets found on left side of slide 5")

        total_score += for_bullet_score
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: AGAINST header with red color (#C0392B) and bold (0.20 points)
    try:
        found_against_header = False
        against_header_score = 0.0

        for shape in text_shapes:
            text = get_all_text_from_shape(shape)
            if 'AGAINST' in text.upper() and 'UNIVERSAL BASIC INCOME' in text.upper():
                # Check it's on the right side (left position > center - some tolerance)
                if shape.left > center_x * 0.8:
                    found_against_header = True
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                color = get_run_color(run)
                                is_bold = run.font.bold is True
                                is_red = color is not None and color.upper() == 'C0392B'
                                size_ok = run.font.size is not None and abs(run.font.size - Pt(22)) < Pt(3)

                                if is_red and is_bold:
                                    against_header_score = 0.20
                                    print(f"PASS: Component 4 — AGAINST header: red={is_red}, bold={is_bold}, size={run.font.size} (0.20 pts)")
                                elif is_red or is_bold:
                                    against_header_score = 0.10
                                    print(f"PARTIAL: Component 4 — AGAINST header: red={is_red}, bold={is_bold} (0.10 pts)")
                                else:
                                    print(f"FAIL: Component 4 — AGAINST header found but color={color}, bold={run.font.bold}")
                                break
                        break
                    break

        if not found_against_header:
            print(f"FAIL: Component 4 — No 'AGAINST Universal Basic Income' header found on right side of slide 5")

        total_score += against_header_score
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: 4 counter-argument bullets on right side (~14pt) (0.15 points)
    try:
        found_against_bullets = False
        against_bullet_score = 0.0

        for shape in text_shapes:
            text = get_all_text_from_shape(shape)
            # Right side text box with multiple paragraphs (not the header)
            if shape.left > center_x * 0.8 and 'AGAINST' not in text.upper().split('\n')[0][:15]:
                para_count = count_nonempty_paragraphs(shape)
                if para_count >= 3:
                    found_against_bullets = True
                    size_ok = False
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size is not None and abs(run.font.size - Pt(14)) < Pt(3):
                                size_ok = True
                                break
                        if size_ok:
                            break

                    if para_count >= 4 and size_ok:
                        against_bullet_score = 0.15
                        print(f"PASS: Component 5 — {para_count} AGAINST bullets found at ~14pt (0.15 pts)")
                    elif para_count >= 4:
                        against_bullet_score = 0.10
                        print(f"PARTIAL: Component 5 — {para_count} AGAINST bullets found but size not 14pt (0.10 pts)")
                    elif para_count >= 2:
                        against_bullet_score = 0.07
                        print(f"PARTIAL: Component 5 — Only {para_count} AGAINST bullets found (0.07 pts)")
                    else:
                        print(f"FAIL: Component 5 — Only {para_count} AGAINST bullets")
                    break

        if not found_against_bullets:
            print(f"FAIL: Component 5 — No counter-argument bullets found on right side of slide 5")

        total_score += against_bullet_score
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: VS circle/oval shape with gold (#F1C40F) fill and bold text (0.15 points)
    try:
        found_vs_shape = False
        vs_score = 0.0

        for shape in auto_shapes:
            text = get_all_text_from_shape(shape)
            if 'VS' in text.upper():
                found_vs_shape = True
                fill_color = get_shape_fill_rgb(shape)
                is_gold_fill = fill_color is not None and fill_color.upper() == 'F1C40F'

                # Check bold text
                is_bold = False
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.bold is True:
                            is_bold = True
                            break

                if is_gold_fill and is_bold:
                    vs_score = 0.15
                    print(f"PASS: Component 6 — VS shape with gold fill ({fill_color}) and bold text (0.15 pts)")
                elif is_gold_fill or is_bold:
                    vs_score = 0.10
                    print(f"PARTIAL: Component 6 — VS shape: gold={is_gold_fill} (fill={fill_color}), bold={is_bold} (0.10 pts)")
                else:
                    vs_score = 0.05
                    print(f"PARTIAL: Component 6 — VS shape found but fill={fill_color}, bold={is_bold} (0.05 pts)")
                break

        # Also check text boxes for VS text if not found in auto shapes
        if not found_vs_shape:
            for shape in text_shapes:
                text = get_all_text_from_shape(shape)
                if text.strip().upper() == 'VS':
                    found_vs_shape = True
                    fill_color = get_shape_fill_rgb(shape)
                    is_gold_fill = fill_color is not None and fill_color.upper() == 'F1C40F'
                    is_bold = False
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.bold is True:
                                is_bold = True
                                break
                    if is_gold_fill and is_bold:
                        vs_score = 0.15
                        print(f"PASS: Component 6 — VS text box with gold fill and bold (0.15 pts)")
                    elif is_gold_fill or is_bold:
                        vs_score = 0.10
                        print(f"PARTIAL: Component 6 — VS text box: gold={is_gold_fill}, bold={is_bold} (0.10 pts)")
                    break

        if not found_vs_shape:
            print(f"FAIL: Component 6 — No VS shape/circle found on slide 5")

        total_score += vs_score
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point — persist app state then verify
def persist_app_state(domain):
    import os, time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
