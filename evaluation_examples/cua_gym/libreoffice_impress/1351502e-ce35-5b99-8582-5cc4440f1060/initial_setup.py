"""
Initial Setup: Create debate presentation with empty slide 5
Task ID: impress_stu_094
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_094'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_content_slide(prs, title_text, body_lines):
    """Add a slide with title and bullet content."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Debate Preparation"
    slide1.placeholders[1].text = "Universal Basic Income — A Comprehensive Analysis"

    # --- Slide 2: Introduction ---
    add_title_content_slide(prs, "Introduction", [
        "Universal Basic Income (UBI) is a government program providing every adult citizen a set amount of money regularly.",
        "The concept has gained traction in recent years amid growing automation concerns.",
        "This presentation examines both sides of the UBI debate to prepare for structured argumentation.",
        "We will analyze economic, social, and ethical dimensions of the policy.",
    ])

    # --- Slide 3: Historical Context ---
    add_title_content_slide(prs, "Historical Context", [
        "Thomas Paine proposed a basic income in 'Agrarian Justice' (1797).",
        "Milton Friedman advocated for a negative income tax in the 1960s.",
        "Alaska's Permanent Fund Dividend has distributed oil revenue since 1982.",
        "Finland ran a two-year UBI pilot program from 2017 to 2018.",
        "Stockton, California tested a guaranteed income program in 2019-2021.",
    ])

    # --- Slide 4: Current Landscape ---
    add_title_content_slide(prs, "Current Landscape", [
        "Over 50 pilot programs have been conducted worldwide since 2020.",
        "AI and automation are projected to displace 85 million jobs by 2025 (World Economic Forum).",
        "COVID-19 stimulus checks served as de facto short-term UBI experiments.",
        "Tech leaders including Elon Musk and Sam Altman have endorsed the concept.",
        "Cost estimates range from $2.8 trillion to $3.8 trillion annually in the United States.",
    ])

    # --- Slide 5: Key Arguments (EMPTY — this is where the agent works) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add only a title text box at the top
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Arguments"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    # Body is intentionally empty — agent must create the split layout

    # --- Slide 6: Discussion Points ---
    add_title_content_slide(prs, "Discussion Points", [
        "How should UBI be funded — through taxation, monetary policy, or budget reallocation?",
        "Would UBI replace existing welfare programs or supplement them?",
        "What is the optimal payment amount to balance incentives and fiscal sustainability?",
        "How do cultural attitudes toward work affect UBI reception across countries?",
    ])

    # --- Slide 7: Conclusion ---
    add_title_content_slide(prs, "Conclusion & Next Steps", [
        "Both sides present compelling evidence rooted in economic theory and pilot data.",
        "The debate ultimately hinges on values: individual freedom vs. collective responsibility.",
        "Next: Each team member will prepare a 3-minute opening statement.",
        "Final debate scheduled for March 28, 2026 — use this deck as reference material.",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
