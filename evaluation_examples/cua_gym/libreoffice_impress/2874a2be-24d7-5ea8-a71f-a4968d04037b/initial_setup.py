"""
Initial Setup: Create photo album source images and blank presentation
Task ID: impress_gf3_048
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from PIL import Image, ImageDraw, ImageFont
import random

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_048'
DESKTOP = f'{WORKDIR}/Desktop'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_sample_photo(filepath, width, height, base_color, label):
    """Create a realistic-looking sample photo with gradient and label."""
    img = Image.new('RGB', (width, height), base_color)
    draw = ImageDraw.Draw(img)

    # Add gradient overlay for visual interest
    r0, g0, b0 = base_color
    for y in range(height):
        factor = y / height
        r = int(r0 * (1 - factor * 0.4))
        g = int(g0 * (1 - factor * 0.3))
        b = int(b0 * (1 - factor * 0.5))
        draw.line([(0, y), (width, y)], fill=(r, g, b))

    # Add some geometric shapes for visual content
    random.seed(hash(label))
    for _ in range(5):
        x1 = random.randint(0, width - 100)
        y1 = random.randint(0, height - 100)
        x2 = x1 + random.randint(50, 150)
        y2 = y1 + random.randint(50, 150)
        shape_color = (
            random.randint(100, 255),
            random.randint(100, 255),
            random.randint(100, 255),
        )
        draw.ellipse([x1, y1, x2, y2], fill=shape_color, outline='white')

    # Add label text
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 36)
    except (IOError, OSError):
        font = ImageFont.load_default()

    bbox = draw.textbbox((0, 0), label, font=font)
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]
    tx = (width - tw) // 2
    ty = (height - th) // 2
    # Shadow
    draw.text((tx + 2, ty + 2), label, fill=(0, 0, 0), font=font)
    draw.text((tx, ty), label, fill=(255, 255, 255), font=font)

    img.save(filepath, 'JPEG', quality=90)
    print(f'Created: {filepath}')


def create_initial():
    os.makedirs(DESKTOP, exist_ok=True)

    # Create 6 varied sample photos with different sizes and colors
    photos = [
        ('photo1.jpg', 1024, 768, (70, 130, 180), 'Mountain Sunrise'),
        ('photo2.jpg', 800, 600, (34, 139, 34), 'Forest Trail'),
        ('photo3.jpg', 1200, 800, (210, 105, 30), 'Autumn Leaves'),
        ('photo4.jpg', 900, 900, (100, 149, 237), 'Ocean Waves'),
        ('photo5.jpg', 1100, 733, (178, 34, 34), 'Desert Sunset'),
        ('photo6.jpg', 960, 720, (75, 0, 130), 'Night Sky'),
    ]

    for fname, w, h, color, label in photos:
        create_sample_photo(os.path.join(DESKTOP, fname), w, h, color, label)

    # Create a blank presentation with 1 default slide
    from pptx import Presentation
    prs = Presentation()
    # python-pptx creates no slides by default; add one blank slide
    blank_layout = prs.slide_layouts[6]  # Blank layout
    prs.slides.add_slide(blank_layout)
    prs.save(OUTPUT)
    print(f'Initial presentation created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
