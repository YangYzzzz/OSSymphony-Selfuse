"""
Reward Script: Photo album-style presentation from 6 images
Task ID: impress_gf3_048
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Slide count >= 6
  Component 2 (0.25): All 6 slides have black background (#000000)
  Component 3 (0.25): Each slide has a picture (image shape present)
  Component 4 (0.25): Each slide has caption text box with correct filename in white text
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_048'


def persist_app_state(domain):
    """Try to save any unsaved changes in LibreOffice."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)
    num_slides = len(slides)
    print(f"INFO: Found {num_slides} slides")

    # Component 1: At least 6 slides (0.25 points)
    # Initial has 1 slide, golden has 6 slides
    try:
        if num_slides >= 6:
            print(f"PASS: Component 1 - Slide count >= 6 (found {num_slides}) (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 - Expected >= 6 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: All 6 photo slides have black background (0.25 points)
    # We check the last 6 slides (in case there's a title slide before them)
    # Each slide with black bg earns 0.25/6 points
    try:
        black_bg_count = 0
        # Check slides that should have photos (at least 6 slides needed)
        check_slides = slides[-6:] if num_slides >= 6 else slides
        for i, slide in enumerate(check_slides):
            fill = slide.background.fill
            bg_color = None
            if fill.type == 1:  # SOLID fill
                try:
                    bg_color = str(fill.fore_color.rgb)
                except:
                    pass
            if bg_color == '000000':
                black_bg_count += 1

        if num_slides >= 6 and black_bg_count == 6:
            print(f"PASS: Component 2 - All 6 photo slides have black background (0.25 pts)")
            total_score += 0.25
        elif black_bg_count > 0 and num_slides >= 6:
            partial = 0.25 * (black_bg_count / 6)
            print(f"PARTIAL: Component 2 - {black_bg_count}/6 slides have black background ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 - {black_bg_count} slides have black background (need 6)")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Each of the 6 photo slides has an image (0.25 points)
    # Each slide with a picture shape earns 0.25/6 points
    try:
        image_count = 0
        check_slides = slides[-6:] if num_slides >= 6 else slides
        for slide in check_slides:
            has_picture = False
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    has_picture = True
                    break
            if has_picture:
                image_count += 1

        if num_slides >= 6 and image_count == 6:
            print(f"PASS: Component 3 - All 6 photo slides contain an image (0.25 pts)")
            total_score += 0.25
        elif image_count > 0 and num_slides >= 6:
            partial = 0.25 * (image_count / 6)
            print(f"PARTIAL: Component 3 - {image_count}/6 slides contain an image ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 - {image_count} slides contain an image (need 6)")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Each slide has a caption text box with correct filename in white (0.25 points)
    # Expected captions: photo1.jpg through photo6.jpg in order
    try:
        expected_captions = [f"photo{i}.jpg" for i in range(1, 7)]
        caption_matches = 0
        check_slides = slides[-6:] if num_slides >= 6 else slides

        for idx, slide in enumerate(check_slides):
            expected = expected_captions[idx] if idx < len(expected_captions) else None
            if expected is None:
                continue

            found_caption = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    full_text = shape.text_frame.text.strip().lower()
                    if expected.lower() in full_text:
                        # Check if text color is white
                        is_white = False
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                try:
                                    rgb = str(run.font.color.rgb)
                                    if rgb.upper() == 'FFFFFF':
                                        is_white = True
                                except:
                                    pass
                        if is_white:
                            found_caption = True
                            break

            if found_caption:
                caption_matches += 1
            else:
                print(f"  DETAIL: Slide {idx+1} - caption '{expected}' with white text not found")

        if num_slides >= 6 and caption_matches == 6:
            print(f"PASS: Component 4 - All 6 captions correct with white text (0.25 pts)")
            total_score += 0.25
        elif caption_matches > 0 and num_slides >= 6:
            partial = 0.25 * (caption_matches / 6)
            print(f"PARTIAL: Component 4 - {caption_matches}/6 captions correct ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 - {caption_matches} captions matched (need 6)")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
