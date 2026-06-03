"""
Reward Script: Architecture Portfolio Presentation
Task ID: impress_wf_053
Domain: libreoffice_impress
Scoring:
  C1 - Slide count == 10 (0.15)
  C2 - Slide 1: title with Portfolio + Alex Chen (0.10)
  C3 - Slide 2: design philosophy with quote (0.10)
  C4 - Slides 3-8: project slides structure (0.25)
  C5 - Slide 9: awards table (0.10)
  C6 - Slide 10: contact slide (0.10)
  C7 - Gold accent color CFB53B (0.10)
  C8 - Uncover transitions on all slides (0.10)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_053'
FILE_NAME = 'Architecture_Portfolio.pptx'

# Check file existence before importing pptx (may not be installed on initial_env)
file_path = os.path.join(WORKDIR, 'Desktop', FILE_NAME)
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
    import sys
    sys.exit(0)

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def get_all_text(slide):
    """Recursively extract all text from a slide (including groups)."""
    texts = []
    def extract(shape):
        if hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            for ri in range(len(table.rows)):
                for ci in range(len(table.columns)):
                    t = table.cell(ri, ci).text.strip()
                    if t:
                        texts.append(t)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                extract(sub)
    for shape in slide.shapes:
        extract(shape)
    return texts


def get_all_colors(slide):
    """Extract all RGB font colors from a slide."""
    colors = set()
    def extract(shape):
        if hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.color.type is not None:
                            colors.add(str(run.font.color.rgb).upper())
                    except Exception:
                        pass
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                extract(sub)
    for shape in slide.shapes:
        extract(shape)
    return colors


def count_tables(slide):
    """Count tables and return list of (rows, cols) tuples."""
    tables = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            t = shape.table
            tables.append((len(t.rows), len(t.columns)))
    return tables


def count_rectangles(slide):
    """Count AUTO_SHAPE rectangles (used as placeholders/dividers)."""
    count = 0
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            count += 1
    return count


def check_transitions(file_path, expected_type='uncover'):
    """Check if all slides have the expected transition type via XML."""
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    slides_with_transition = 0
    total_slides = 0
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            slide_files = [n for n in zf.namelist() if n.startswith('ppt/slides/slide') and n.endswith('.xml')]
            total_slides = len(slide_files)
            for sf in slide_files:
                with zf.open(sf) as f:
                    root = ET.parse(f).getroot()
                    tr = root.find('.//p:transition', ns)
                    if tr is not None:
                        # Check for the expected transition child element
                        for child in tr:
                            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                            if tag.lower() == expected_type.lower():
                                slides_with_transition += 1
                                break
    except Exception as e:
        print(f"ERROR: Transition check failed: {e}")
        return 0, 0
    return slides_with_transition, total_slides


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Exactly 10 slides (0.15 points)
    try:
        if num_slides == 10:
            print(f"PASS: Component 1 -- Slide count is 10 (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 -- Expected 10 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Slide 1 title with "Portfolio" and "Alex Chen" (0.10 points)
    try:
        if num_slides >= 1:
            slide1_texts = get_all_text(prs.slides[0])
            full_text = ' '.join(slide1_texts).lower()
            has_portfolio = 'portfolio' in full_text
            has_alex_chen = 'alex chen' in full_text
            if has_portfolio and has_alex_chen:
                print(f"PASS: Component 2 -- Slide 1 has 'Portfolio' and 'Alex Chen' (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 2 -- Slide 1 missing: portfolio={has_portfolio}, alex_chen={has_alex_chen}. Texts: {slide1_texts[:3]}")
        else:
            print(f"FAIL: Component 2 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Slide 2 design philosophy with quote (0.10 points)
    try:
        if num_slides >= 2:
            slide2_texts = get_all_text(prs.slides[1])
            full_text = ' '.join(slide2_texts).lower()
            # Should have design philosophy heading and a quote-like text
            has_philosophy = 'design philosophy' in full_text or 'philosophy' in full_text
            # Check for quote (italic text or quotation marks or long text)
            quote_found = any(
                run.font.italic and len(run.text.strip()) > 20
                or '"' in (run.text or '')
                or len(para.text.strip()) > 50
                for shape in prs.slides[1].shapes if hasattr(shape, 'text_frame')
                for para in shape.text_frame.paragraphs
                for run in para.runs
            )
            if has_philosophy and quote_found:
                print(f"PASS: Component 3 -- Slide 2 has design philosophy with quote (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 3 -- philosophy={has_philosophy}, quote={quote_found}")
        else:
            print(f"FAIL: Component 3 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Slides 3-8 project structure (0.25 points)
    # Each project slide (6 slides) should have: project number, header text, rectangles (placeholders), mini table
    try:
        if num_slides >= 8:
            valid_project_slides = 0
            for idx in range(2, 8):  # slides 3-8 (0-indexed: 2-7)
                slide = prs.slides[idx]
                slide_texts = get_all_text(slide)
                full_text = ' '.join(slide_texts).lower()

                # Check for project number (01-06)
                expected_num = f"{idx - 1:02d}"
                has_number = any(expected_num in t for t in slide_texts)

                # Check for header with project info (name | year | location pattern)
                has_header = any('|' in t or ('20' in t and len(t) > 15) for t in slide_texts)

                # Check for rectangles (image placeholder + grid = multiple rectangles)
                rect_count = count_rectangles(slide)
                has_placeholders = rect_count >= 3  # at least large rect + grid lines

                # Check for mini table with 3 rows
                tables = count_tables(slide)
                has_table = any(rows >= 3 for rows, cols in tables)

                if has_number and has_header and has_placeholders and has_table:
                    valid_project_slides += 1
                else:
                    print(f"  Slide {idx+1}: number={has_number}, header={has_header}, placeholders={has_placeholders}(rects={rect_count}), table={has_table}(tables={tables})")

            # Award proportional credit: 0.25 * (valid/6)
            slide_score = 0.25 * (valid_project_slides / 6.0)
            if valid_project_slides == 6:
                print(f"PASS: Component 4 -- All 6 project slides valid (0.25 pts)")
                total_score += slide_score
            elif valid_project_slides > 0:
                print(f"PARTIAL: Component 4 -- {valid_project_slides}/6 project slides valid ({slide_score:.3f} pts)")
                total_score += slide_score
            else:
                print(f"FAIL: Component 4 -- 0/6 project slides valid")
        else:
            print(f"FAIL: Component 4 -- Not enough slides for projects (need 8, have {num_slides})")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Slide 9 awards table with Year/Award/Project (0.10 points)
    try:
        if num_slides >= 9:
            slide9 = prs.slides[8]
            slide9_texts = get_all_text(slide9)
            full_text = ' '.join(slide9_texts).lower()
            has_awards_heading = 'award' in full_text
            tables = count_tables(slide9)
            # Awards table should have 3 columns and multiple rows
            has_awards_table = any(cols >= 3 and rows >= 3 for rows, cols in tables)

            # Check table header content
            correct_headers = any(
                len(shape.table.columns) >= 3
                and 'year' in shape.table.cell(0, 0).text.strip().lower()
                and 'award' in shape.table.cell(0, 1).text.strip().lower()
                and 'project' in shape.table.cell(0, 2).text.strip().lower()
                for shape in slide9.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE
            )

            if has_awards_heading and has_awards_table and correct_headers:
                print(f"PASS: Component 5 -- Slide 9 has awards table with correct headers (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 5 -- heading={has_awards_heading}, table={has_awards_table}, headers={correct_headers}")
        else:
            print(f"FAIL: Component 5 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # Component 6: Slide 10 contact slide (0.10 points)
    try:
        if num_slides >= 10:
            slide10 = prs.slides[9]
            slide10_texts = get_all_text(slide10)
            full_text = ' '.join(slide10_texts).lower()
            # Contact slide should have contact info keywords
            has_contact = ('contact' in full_text or 'get in touch' in full_text or
                          'email' in full_text or 'phone' in full_text)
            has_details = ('email' in full_text or '@' in full_text) and ('phone' in full_text or 'studio' in full_text)

            if has_contact and has_details:
                print(f"PASS: Component 6 -- Slide 10 is a contact slide (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 6 -- contact={has_contact}, details={has_details}")
        else:
            print(f"FAIL: Component 6 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    # Component 7: Gold accent color CFB53B present (0.10 points)
    try:
        gold_found_on_slides = 0
        for idx, slide in enumerate(prs.slides):
            colors = get_all_colors(slide)
            if 'CFB53B' in colors:
                gold_found_on_slides += 1

        # Gold accent should appear on multiple slides (project numbers, headings, dividers)
        if gold_found_on_slides >= 3:
            print(f"PASS: Component 7 -- Gold accent #CFB53B found on {gold_found_on_slides} slides (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 7 -- Gold accent #CFB53B found on only {gold_found_on_slides} slides (need >= 3)")
    except Exception as e:
        print(f"ERROR: Component 7 -- {e}")

    # Component 8: Uncover transitions on all slides (0.10 points)
    try:
        slides_with_uncover, total_slide_count = check_transitions(file_path, 'uncover')
        if total_slide_count > 0 and slides_with_uncover == total_slide_count:
            print(f"PASS: Component 8 -- Uncover transitions on all {total_slide_count} slides (0.10 pts)")
            total_score += 0.10
        elif slides_with_uncover > 0:
            partial = 0.10 * (slides_with_uncover / max(total_slide_count, 1))
            print(f"PARTIAL: Component 8 -- Uncover on {slides_with_uncover}/{total_slide_count} slides ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 8 -- No Uncover transitions found ({slides_with_uncover}/{total_slide_count})")
    except Exception as e:
        print(f"ERROR: Component 8 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point — file existence already checked at top of script
verify_task(file_path)
