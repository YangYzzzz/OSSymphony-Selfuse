"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 218 I need to trim “Picture 1” into a perfect square (same height and width) and then align it exactly in the center of the slide in LibreOffice Impress. How can I do that?
Generated: 2025-09-10 20:08:47
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
from pptx import Presentation

# -----------------------------------------------------------
# Reward Script: Verify picture on slide 218 is trimmed to a
#                perfect square and perfectly centred
# -----------------------------------------------------------

def verify_slide218_center_square(file_path: str) -> float:
    """Verify that on slide 218 the shape named 'Picture 1'
    1. Has equal width and height  (square)
    2. Is centred both horizontally and vertically on the slide

    Returns a progressive score between 0.0 and 1.0.
    """

    score = 0.0       # progressive score accumulator
    max_score = 1.0   # maximum score that can be awarded

    # --------------- 0. Prerequisite: file must load ----------------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        return 0.0

    # --------------- 1. Retrieve slide 218 (index 217) -------------
    slide_index = 217  # zero-based index for slide 218
    if len(prs.slides) <= slide_index:
        print(f"✗ Presentation has only {len(prs.slides)} slides; slide 218 missing")
        return 0.0

    slide = prs.slides[slide_index]
    slide_w, slide_h = prs.slide_width, prs.slide_height
    print(f"Slide dimensions (EMU): width={slide_w}, height={slide_h}")

    # --------------- 2. Locate the target picture -------------------
    target = None
    for shp in slide.shapes:
        if shp.shape_type == 13:  # 13 == PICTURE
            if shp.name.strip().lower() == "picture 1":
                target = shp
                break

    if target is None:
        print("✗ Picture named 'Picture 1' not found on slide 218")
        return 0.0

    print(f"✓ Found picture: {target.name}")

    # --------------- 3. Check for perfect square --------------------
    width, height = target.width, target.height
    if width == 0 or height == 0:
        print("✗ Picture has zero width or height — cannot verify square")
    else:
        diff = abs(width - height)
        tolerance = max(width, height) * 0.01  # 1 % tolerance
        if diff <= tolerance:
            print(f"✓ Picture is square → diff {diff} ≤ tolerance {tolerance}")
            score += 0.4  # 40 % for correct trimming
        else:
            print(f"✗ Picture not square → width {width}, height {height}")

    # --------------- 4. Check centre alignment ---------------------
    centre_x = target.left + width / 2
    centre_y = target.top + height / 2
    slide_centre_x = slide_w / 2
    slide_centre_y = slide_h / 2

    tol_x = max(slide_w * 0.01, 2000)  # 1 % of slide width or 2000 EMU
    tol_y = max(slide_h * 0.01, 2000)

    # Horizontal centre
    if abs(centre_x - slide_centre_x) <= tol_x:
        print("✓ Picture horizontally centred")
        score += 0.3  # 30 % for horizontal centring
    else:
        print(f"✗ Horizontal centring off by {abs(centre_x - slide_centre_x)} EMU (tol {tol_x})")

    # Vertical centre
    if abs(centre_y - slide_centre_y) <= tol_y:
        print("✓ Picture vertically centred")
        score += 0.3  # 30 % for vertical centring
    else:
        print(f"✗ Vertical centring off by {abs(centre_y - slide_centre_y)} EMU (tol {tol_y})")

    # --------------- 5. Final score & return -----------------------
    final_score = min(score, max_score)
    print(f"Computed score: {final_score}")
    return final_score


# ---------------------- Script Entrypoint ---------------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/on_slide_218_i_need_to_trim_picture_1_into_a_perfect_square_same_height_and_width_and_then_align_it__golden.pptx"
    reward = verify_slide218_center_square(FILE_PATH)
    print(f"REWARD: {reward}")

