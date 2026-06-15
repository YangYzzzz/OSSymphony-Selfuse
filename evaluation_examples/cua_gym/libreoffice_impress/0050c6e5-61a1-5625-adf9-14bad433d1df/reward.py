"""
Reward Script: Add a rectangular callout box on slide 2 with text 'Key Concept',
light yellow fill (#FFFDE7), and 2pt dark orange (#E65100) border.
Task ID: impress_teach_019
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.4): Rectangle shape on slide 2 with text 'Key Concept'
  - Component 2 (0.3): Fill color is #FFFDE7
  - Component 3 (0.3): Border is 2pt with color #E65100
"""

import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_019'


def persist_app_state(domain: str):
    """Try to save any unsaved changes in LibreOffice."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def find_key_concept_rectangle(slide):
    """
    Search all shapes on the slide for a rectangle-type AUTO_SHAPE
    containing text 'Key Concept'. Returns the shape or None.
    """
    for shape in slide.shapes:
        # Accept AUTO_SHAPE (rectangles, rounded rects, etc.)
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if hasattr(shape, 'text') and 'key concept' in shape.text.strip().lower():
                return shape
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 2")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]  # 0-indexed, slide 2

    # Component 1: Rectangle shape on slide 2 with text 'Key Concept' (0.4 points)
    rect_shape = None
    try:
        rect_shape = find_key_concept_rectangle(slide2)
        if rect_shape is not None:
            print(f"PASS: Component 1 — Rectangle with text '{rect_shape.text}' found on slide 2 (0.4 pts)")
            total_score += 0.4
        else:
            print("FAIL: Component 1 — No rectangle shape with text 'Key Concept' found on slide 2")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Fill color is #FFFDE7 (0.3 points)
    if rect_shape is not None:
        try:
            fill = rect_shape.fill
            if fill.type is not None and fill.type == 1:  # SOLID fill
                actual_color = str(fill.fore_color.rgb).upper()
                if actual_color == 'FFFDE7':
                    print(f"PASS: Component 2 — Fill color is #{actual_color} (0.3 pts)")
                    total_score += 0.3
                else:
                    print(f"FAIL: Component 2 — Fill color is #{actual_color}, expected #FFFDE7")
            else:
                print(f"FAIL: Component 2 — Fill type is {fill.type}, expected SOLID (1)")
        except Exception as e:
            print(f"ERROR: Component 2 — {e}")
    else:
        print("SKIP: Component 2 — No rectangle found, cannot check fill color")

    # Component 3: Border 2pt with color #E65100 (0.3 points)
    if rect_shape is not None:
        try:
            line = rect_shape.line
            line_width = line.width
            # 2pt = 25400 EMU; allow some tolerance (1.5pt to 2.5pt)
            width_ok = line_width is not None and 19050 <= line_width <= 31750

            color_ok = False
            try:
                if line.color.type is not None:
                    actual_line_color = str(line.color.rgb).upper()
                    color_ok = (actual_line_color == 'E65100')
                    if not color_ok:
                        print(f"FAIL: Component 3 — Border color is #{actual_line_color}, expected #E65100")
                else:
                    print("FAIL: Component 3 — Border color type is None")
            except Exception as e:
                print(f"FAIL: Component 3 — Cannot read border color: {e}")

            if width_ok and color_ok:
                print(f"PASS: Component 3 — Border is {line_width/12700:.1f}pt with color #E65100 (0.3 pts)")
                total_score += 0.3
            elif width_ok and not color_ok:
                pass  # already printed fail above
            elif not width_ok:
                actual_pt = line_width / 12700 if line_width else 0
                print(f"FAIL: Component 3 — Border width is {actual_pt:.1f}pt, expected ~2.0pt")
        except Exception as e:
            print(f"ERROR: Component 3 — {e}")
    else:
        print("SKIP: Component 3 — No rectangle found, cannot check border")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
