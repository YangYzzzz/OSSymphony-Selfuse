"""
Initial Setup: Add a rectangular callout box on slide 2
Task ID: impress_teach_019
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_019'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Micro Economics"
    slide1.placeholders[1].text = "Principles of Supply, Demand, and Market Equilibrium"

    # --- Slide 2: Supply and Demand (lecture content, NO callout box) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Title text box
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Supply and Demand Fundamentals"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    # Content text box
    content_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(8), Inches(5.5))
    tf2 = content_box.text_frame
    tf2.word_wrap = True

    paragraphs_text = [
        "The law of demand states that, all else being equal, as the price of a product increases, the quantity demanded falls; conversely, as the price falls, quantity demanded increases.",
        "The law of supply indicates that a higher price leads to a higher quantity supplied, while a lower price leads to a lower quantity supplied, holding all other factors constant.",
        "Market equilibrium occurs at the price where quantity supplied equals quantity demanded. At this intersection, there is no tendency for the price to change.",
        "Consumer surplus represents the difference between what consumers are willing to pay and what they actually pay. Producer surplus is the difference between the market price and the minimum price producers would accept.",
        "Shifts in supply or demand curves — caused by changes in income, preferences, technology, or input costs — lead to new equilibrium prices and quantities.",
    ]

    for i, text in enumerate(paragraphs_text):
        if i == 0:
            para = tf2.paragraphs[0]
        else:
            para = tf2.add_paragraph()
        para.text = text
        para.space_after = Pt(10)
        for r in para.runs:
            r.font.name = "Arial"
            r.font.size = Pt(14)
            r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 3: Elasticity ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Price Elasticity of Demand"
    run3 = p3.runs[0]
    run3.font.name = "Arial"
    run3.font.size = Pt(28)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    body3 = slide3.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(8), Inches(5.5))
    tf3b = body3.text_frame
    tf3b.word_wrap = True
    elast_text = [
        "Elasticity measures how responsive quantity demanded is to a change in price. It is calculated as the percentage change in quantity divided by the percentage change in price.",
        "When elasticity > 1, demand is elastic — consumers are highly sensitive to price changes. Common examples include luxury goods and products with many substitutes.",
        "When elasticity < 1, demand is inelastic — quantity demanded changes little despite price fluctuations. Necessities such as insulin or gasoline often fall into this category.",
        "Perfectly inelastic demand (elasticity = 0) means quantity demanded does not change at all regardless of price, while perfectly elastic demand (elasticity = infinity) means any price increase causes demand to drop to zero.",
    ]
    for i, text in enumerate(elast_text):
        if i == 0:
            para = tf3b.paragraphs[0]
        else:
            para = tf3b.add_paragraph()
        para.text = text
        para.space_after = Pt(10)
        for r in para.runs:
            r.font.name = "Arial"
            r.font.size = Pt(14)
            r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 4: Market Structures ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Market Structures Overview"
    run4 = p4.runs[0]
    run4.font.name = "Arial"
    run4.font.size = Pt(28)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    body4 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(8), Inches(5.5))
    tf4b = body4.text_frame
    tf4b.word_wrap = True
    mkt_text = [
        "Perfect Competition: Many firms sell identical products; no single firm can influence market price. Agricultural commodities are a classic example.",
        "Monopoly: A single firm dominates the entire market, often due to barriers to entry such as patents, resource control, or economies of scale.",
        "Oligopoly: A small number of large firms control most of the market share. Strategic interaction among firms is a defining characteristic, often analyzed through game theory.",
        "Monopolistic Competition: Many firms sell differentiated products. Firms have some price-setting power due to brand loyalty and product features, but face competition from close substitutes.",
    ]
    for i, text in enumerate(mkt_text):
        if i == 0:
            para = tf4b.paragraphs[0]
        else:
            para = tf4b.add_paragraph()
        para.text = text
        para.space_after = Pt(10)
        for r in para.runs:
            r.font.name = "Arial"
            r.font.size = Pt(14)
            r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 5: Summary ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf5 = txBox5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Key Takeaways"
    run5 = p5.runs[0]
    run5.font.name = "Arial"
    run5.font.size = Pt(28)
    run5.font.bold = True
    run5.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    body5 = slide5.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(8), Inches(5.5))
    tf5b = body5.text_frame
    tf5b.word_wrap = True
    summary_text = [
        "Supply and demand interact to determine equilibrium price and quantity in competitive markets.",
        "Elasticity quantifies how sensitive consumers and producers are to price changes, guiding policy and business strategy.",
        "Different market structures — from perfect competition to monopoly — shape firm behavior, pricing power, and consumer welfare.",
        "Understanding these microeconomic foundations is essential for analyzing real-world markets and making informed economic decisions.",
    ]
    for i, text in enumerate(summary_text):
        if i == 0:
            para = tf5b.paragraphs[0]
        else:
            para = tf5b.add_paragraph()
        para.text = text
        para.space_after = Pt(10)
        for r in para.runs:
            r.font.name = "Arial"
            r.font.size = Pt(14)
            r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
