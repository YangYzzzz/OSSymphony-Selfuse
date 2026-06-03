"""
Initial Setup: Create a presentation with 8 slides for Department Review.
Task ID: impress_stu_069
Domain: libreoffice_impress
Slide 6 ('Publication Trends') is empty — no chart.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_069'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_body_slide(prs, layout_idx, title_text, body_lines):
    """Add a slide with title and bullet-point body text."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
    return slide


def add_title_only_slide(prs, title_text):
    """Add a slide with only a title (layout 5 = Title Only)."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title_text
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Department Review 2025"
    slide1.placeholders[1].text = "Computer Science Research Division\nAnnual Performance Summary"

    # --- Slide 2: Team Overview ---
    add_title_body_slide(prs, 1, "Team Overview", [
        "Current headcount: 42 full-time researchers, 18 postdocs",
        "New hires in 2024: Dr. Aisha Patel (NLP), Dr. Liam O'Brien (Systems)",
        "3 promotions to Senior Researcher in Q4 2024",
        "Average tenure: 4.2 years across all levels",
        "Cross-functional collaborations with 6 partner institutions",
    ])

    # --- Slide 3: Budget Summary ---
    add_title_body_slide(prs, 1, "Budget Summary", [
        "Total FY2025 budget: $12.4M (up 8% from FY2024)",
        "Personnel costs: $7.8M (63% of total)",
        "Equipment & computing: $2.1M — new GPU cluster approved",
        "Travel & conferences: $0.9M — ICML, NeurIPS, CVPR attendance",
        "Overhead & facilities: $1.6M",
    ])

    # --- Slide 4: Project Milestones ---
    add_title_body_slide(prs, 1, "Project Milestones", [
        "Project Aurora (LLM fine-tuning): Phase 2 complete, demo in March",
        "Project Beacon (autonomous agents): Filed 2 patents in Q3",
        "Project Catalyst (drug discovery ML): Partnership with PharmaCorp signed",
        "Project Delta (edge inference): Achieved 40ms latency on mobile SoC",
        "3 projects transitioned to production in FY2024",
    ])

    # --- Slide 5: Staffing & Recruitment ---
    add_title_body_slide(prs, 1, "Staffing & Recruitment", [
        "Open positions: 5 research scientists, 2 engineering leads",
        "Pipeline: 120 applicants screened, 28 on-site interviews scheduled",
        "Diversity initiative: 45% of new hires from underrepresented groups",
        "Intern program: 12 summer interns from top-10 CS programs",
        "Retention rate: 91% (industry average: 82%)",
    ])

    # --- Slide 6: Publication Trends (EMPTY — no chart) ---
    slide6 = add_title_only_slide(prs, "Publication Trends")
    # Intentionally left empty — the task asks the agent to add a chart here

    # --- Slide 7: Goals for Next Year ---
    add_title_body_slide(prs, 1, "Goals for Next Year", [
        "Increase publication output by 15% (target: 127 papers)",
        "Secure $3M in external grant funding (NSF, DARPA)",
        "Launch 2 new cross-departmental research initiatives",
        "Expand GPU cluster capacity by 50% for LLM training",
        "Host inaugural departmental research symposium in September",
    ])

    # --- Slide 8: Q&A / Discussion ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[0])
    slide8.shapes.title.text = "Q&A / Discussion"
    slide8.placeholders[1].text = "Thank you for your attention.\nQuestions and feedback welcome."

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
