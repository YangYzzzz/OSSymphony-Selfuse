"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 262 is the odd one out—everything in Content Text Box 1 is left-aligned while the rest of the deck is centered. Can you center-align all the text in that specific placeholder so it matches the other slides?
Generated: 2025-09-10 19:25:25
Status: success
Model: azure-o3
Total Steps: 2
"""

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
import os

FILE_PATH = "/home/user/slide_262_is_the_odd_one_outeverything_in_content_text_box_1_is_left_aligned_while_the_rest_of_the_d_golden.pptx"

def verify_center_alignment_slide_262(file_path: str) -> float:
    """Verify that all paragraphs in Content Placeholder 1 on slide 262 are
    center-aligned. Progressive scoring is applied based on:
        • File/slide accessibility (0.0 pts – prerequisite, no score)
        • Having at least 262 slides (0.1 pts)
        • Locating the correct content placeholder (0.2 pts)
        • Having at least one centered paragraph (0.3 pts)
        • All paragraphs centered (0.4 pts)
       Maximum score = 1.0
    """
    print(f"Loading presentation: {file_path}")
    if not os.path.exists(file_path):
        print("✗ File not found")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error opening presentation: {e}")
        return 0.0

    total_score = 0.0
    max_score = 1.0

    # Requirement 1: Ensure deck has at least 262 slides (index 261 exists)
    if len(prs.slides) >= 262:
        print("✓ Slide count sufficient (0.1 points)")
        total_score += 0.1
    else:
        print("✗ Deck contains fewer than 262 slides – cannot verify further")
        return total_score  # Early exit, nothing else can be verified

    # Access slide 262 (index 261)
    target_slide = prs.slides[261]
    print("Analysing slide 262 …")

    # Locate the primary content placeholder (OBJECT or BODY)
    content_shape = None
    for shape in target_slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type in (7, 2):  # OBJECT or BODY
            content_shape = shape
            break

    if content_shape is None:
        print("✗ Content placeholder not found (slide may be incorrectly formatted)")
        return total_score

    print(f"✓ Found content placeholder '{getattr(content_shape, 'name', '')}' (0.2 points)")
    total_score += 0.2

    # Verify paragraph alignment inside the placeholder
    tf = content_shape.text_frame
    if tf is None or len(tf.paragraphs) == 0:
        print("✗ No text found in content placeholder – cannot verify alignment")
        return total_score

    center_count = 0
    for idx, para in enumerate(tf.paragraphs):
        alignment = para.alignment
        print(f"  Paragraph {idx}: alignment = {alignment}")
        if alignment == PP_ALIGN.CENTER:
            center_count += 1

    # Award points based on how many paragraphs are centered
    if center_count > 0:
        print("✓ At least one paragraph is centered (+0.3 points)")
        total_score += 0.3
    else:
        print("✗ No centered paragraphs detected")
        return total_score

    if center_count == len(tf.paragraphs):
        print(f"✓ All {center_count} paragraphs are centered (+0.4 points)")
        total_score += 0.4
    else:
        print(f"✗ Only {center_count}/{len(tf.paragraphs)} paragraphs are centered")

    final_score = min(total_score, max_score)
    print(f"Final score: {final_score}")
    return final_score


if __name__ == "__main__":
    reward = verify_center_alignment_slide_262(FILE_PATH)
    print(f"REWARD: {reward}")
