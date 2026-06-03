"""
FINAL REWARD SCRIPT - SUCCESS
Task: Switch pages 1–2 to Roman numerals (i, ii) and pages 3+ to Arabic starting at 1.
Generated: 2025-10-17 17:27:10
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation

# -------------------------------------------------------------
# Reward Script for LibreOffice Impress Task
# Task: "Switch pages 1–2 to Roman numerals (i, ii) and pages 3+ to Arabic starting at 1."
# -------------------------------------------------------------
# Verification Strategy
# 1. Load the PPTX presentation (no points for merely loading).
# 2. Inspect every slide and gather all visible text strings.
# 3. For slides 1 & 2  ➜ verify presence of the EXACT roman numerals:  "i" and "ii" (case-insensitive).
# 4. For slides 3+   ➜ verify presence of the EXACT arabic numbers starting at "1".
# 5. Progressive scoring:
#    • Roman numerals (first two slides) contribute 0.4 of the total score.
#    • Arabic numbering (remaining slides) contributes 0.6 of the total score.
#    • Scores are proportional to how many of the required numbers are correct.
# 6. Return a final float in [0.0, 1.0] and print it as  "REWARD: X.X".
# -------------------------------------------------------------

def _roman_numeral(n: int) -> str:
    """Return lower-case roman numeral for 1 <= n <= 20 (covers usual cases)."""
    romans = [
        "i", "ii", "iii", "iv", "v", "vi", "vii", "viii", "ix", "x",
        "xi", "xii", "xiii", "xiv", "xv", "xvi", "xvii", "xviii", "xix", "xx",
    ]
    return romans[n - 1] if 1 <= n <= len(romans) else ""


def _collect_slide_texts(slide):
    """Return a list of all non-empty text strings found on the slide."""
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            txt = shape.text.strip()
            if txt:
                texts.append(txt)
    return texts


def verify_page_numbering(file_path: str) -> float:
    print(f"Checking presentation: {file_path}\n")

    # ---------- Safety checks ----------
    if not os.path.exists(file_path):
        print("✗ File not found!")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error opening presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    slide_count = len(prs.slides)
    if slide_count == 0:
        print("✗ Presentation contains no slides")
        print("REWARD: 0.0")
        return 0.0

    print(f"Total slides detected: {slide_count}\n")

    # ---------- Verification ----------
    roman_required = min(2, slide_count)          # Slides 1–2 (if they exist)
    arabic_required = max(0, slide_count - 2)     # Slides 3+

    roman_correct = 0
    arabic_correct = 0

    for idx, slide in enumerate(prs.slides, start=1):
        texts = _collect_slide_texts(slide)
        lower_texts = [t.lower() for t in texts]

        if idx <= 2:  # Roman numerals expected
            expected = _roman_numeral(idx)
            if expected and expected in lower_texts:
                roman_correct += 1
                print(f"✓ Slide {idx}: contains correct Roman numeral '{expected}'")
            else:
                print(f"✗ Slide {idx}: expected Roman numeral '{expected}' — NOT found")
        else:         # Arabic numbers expected starting at 1
            expected = str(idx - 2)
            if expected in texts:  # Case-sensitive numeric match
                arabic_correct += 1
                print(f"✓ Slide {idx}: contains correct Arabic number '{expected}'")
            else:
                print(f"✗ Slide {idx}: expected Arabic number '{expected}' — NOT found")

    # ---------- Scoring ----------
    score = 0.0
    # Weight: 0.4 for roman section, 0.6 for arabic section
    if roman_required:
        score += 0.4 * (roman_correct / roman_required)
    if arabic_required:
        score += 0.6 * (arabic_correct / arabic_required)

    final_score = round(min(score, 1.0), 2)  # Rounded for neatness; cap at 1.0

    print("\nVerification complete.")
    print(f"Roman numerals  : {roman_correct}/{roman_required}")
    print(f"Arabic numbers  : {arabic_correct}/{arabic_required}")
    print(f"Total score     : {final_score}")
    print(f"REWARD: {final_score}")

    return final_score


# -------------------------------
# Execute verification (MANDATORY)
# -------------------------------
if __name__ == "__main__":
    FILE = "/home/user/switch_pages_12_to_roman_numerals_i_ii_and_pages_3_to_arabic_starting_at_1.pptx"
    verify_page_numbering(FILE)

