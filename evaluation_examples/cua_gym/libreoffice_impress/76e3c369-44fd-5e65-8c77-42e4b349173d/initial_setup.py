"""
Initial Setup: Remove all animations from slide 5 without affecting other slides
Task ID: impress_gf3_044
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import copy
import zipfile
import shutil
from io import BytesIO
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_044'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

# XML namespaces for animation
NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}

for prefix, uri in NSMAP.items():
    ET.register_namespace(prefix, uri)
# Also register common namespaces that appear in pptx
ET.register_namespace('p14', 'http://schemas.microsoft.com/office/powerpoint/2010/main')
ET.register_namespace('mc', 'http://schemas.openxmlformats.org/markup-compatibility/2006')
ET.register_namespace('r', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships')


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Mixed Deck Presentation"
    slide1.placeholders[1].text = "Q1 2025 Strategic Overview"

    # --- Slide 2: Company Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Company Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Acme Corporation was founded in 2010 and has grown to over 500 employees."
    p = body2.add_paragraph()
    p.text = "Headquarters: San Francisco, CA"
    p = body2.add_paragraph()
    p.text = "Annual Revenue: $45.2M (2024)"
    p = body2.add_paragraph()
    p.text = "Key Markets: North America, Europe, Asia-Pacific"

    # --- Slide 3: Revenue Analysis (will have animations) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    title3.text_frame.paragraphs[0].text = "Revenue Analysis"
    title3.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    title3.text_frame.paragraphs[0].runs[0].font.bold = True
    title3.name = "RevTitle"

    box3a = slide3.shapes.add_textbox(Inches(1), Inches(1.8), Inches(5), Inches(1))
    box3a.text_frame.paragraphs[0].text = "North America: $28.5M (+12% YoY)"
    box3a.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    box3a.name = "RevNA"

    box3b = slide3.shapes.add_textbox(Inches(1), Inches(3.0), Inches(5), Inches(1))
    box3b.text_frame.paragraphs[0].text = "Europe: $11.7M (+8% YoY)"
    box3b.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    box3b.name = "RevEU"

    # --- Slide 4: Team Structure ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Team Structure"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Engineering: 180 employees across 12 teams"
    body4.add_paragraph().text = "Product: 45 employees, 3 product lines"
    body4.add_paragraph().text = "Sales & Marketing: 95 employees, 5 regions"
    body4.add_paragraph().text = "Operations: 120 employees, 4 divisions"
    body4.add_paragraph().text = "Executive: 15 members including board advisors"

    # --- Slide 5: Product Roadmap (will have 6 animations - target slide) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10), Inches(1))
    title5.text_frame.paragraphs[0].text = "Product Roadmap 2025"
    title5.text_frame.paragraphs[0].runs[0].font.size = Pt(36)
    title5.text_frame.paragraphs[0].runs[0].font.bold = True
    title5.name = "RoadmapTitle"

    box5a = slide5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(1))
    box5a.text_frame.paragraphs[0].text = "Phase 1: Platform Redesign (Jan-Mar)"
    box5a.text_frame.paragraphs[0].runs[0].font.size = Pt(20)
    box5a.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
    box5a.name = "Phase1"

    box5b = slide5.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(5.5), Inches(1))
    box5b.text_frame.paragraphs[0].text = "Phase 2: AI Integration (Apr-Jun)"
    box5b.text_frame.paragraphs[0].runs[0].font.size = Pt(20)
    box5b.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x2E, 0x75, 0xB6)
    box5b.name = "Phase2"

    box5c = slide5.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(5.5), Inches(1))
    box5c.text_frame.paragraphs[0].text = "Phase 3: Global Expansion (Jul-Sep)"
    box5c.text_frame.paragraphs[0].runs[0].font.size = Pt(20)
    box5c.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x54, 0x8B, 0x54)
    box5c.name = "Phase3"

    box5d = slide5.shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(5.5), Inches(1))
    box5d.text_frame.paragraphs[0].text = "Budget Allocation: $12.5M"
    box5d.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    box5d.name = "Budget"

    box5e = slide5.shapes.add_textbox(Inches(6.5), Inches(2.8), Inches(5.5), Inches(1))
    box5e.text_frame.paragraphs[0].text = "Expected ROI: 340% within 18 months"
    box5e.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    box5e.name = "ROI"

    box5f = slide5.shapes.add_textbox(Inches(6.5), Inches(4.1), Inches(5.5), Inches(1))
    box5f.text_frame.paragraphs[0].text = "Key Metric: 2M active users by Q4"
    box5f.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    box5f.name = "Metric"

    # --- Slide 6: Market Analysis ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Market Analysis"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Total Addressable Market (TAM): $8.2B"
    body6.add_paragraph().text = "Serviceable Available Market (SAM): $2.1B"
    body6.add_paragraph().text = "Serviceable Obtainable Market (SOM): $420M"
    body6.add_paragraph().text = "Current market share: 2.3%"

    # --- Slide 7: Customer Success Stories (will have animations) ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    title7.text_frame.paragraphs[0].text = "Customer Success Stories"
    title7.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    title7.text_frame.paragraphs[0].runs[0].font.bold = True
    title7.name = "CustTitle"

    box7a = slide7.shapes.add_textbox(Inches(1), Inches(1.8), Inches(5), Inches(2))
    tf7a = box7a.text_frame
    tf7a.word_wrap = True
    tf7a.paragraphs[0].text = "TechCorp Inc. reduced onboarding time by 65% using our platform."
    tf7a.paragraphs[0].runs[0].font.size = Pt(16)
    box7a.name = "Story1"

    box7b = slide7.shapes.add_textbox(Inches(1), Inches(4.0), Inches(5), Inches(2))
    tf7b = box7b.text_frame
    tf7b.word_wrap = True
    tf7b.paragraphs[0].text = "GlobalRetail saved $3.2M annually through automated inventory management."
    tf7b.paragraphs[0].runs[0].font.size = Pt(16)
    box7b.name = "Story2"

    # --- Slide 8: Financial Summary ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Financial Summary"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Revenue Growth: 18% CAGR over 3 years"
    body8.add_paragraph().text = "Gross Margin: 72%"
    body8.add_paragraph().text = "Operating Expenses: $32.1M"
    body8.add_paragraph().text = "Net Income: $4.8M"
    body8.add_paragraph().text = "Cash on Hand: $22.3M"

    # --- Slide 9: Future Initiatives (will have animations) ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    title9.text_frame.paragraphs[0].text = "Future Initiatives"
    title9.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    title9.text_frame.paragraphs[0].runs[0].font.bold = True
    title9.name = "FutureTitle"

    box9a = slide9.shapes.add_textbox(Inches(1), Inches(1.8), Inches(5), Inches(1.5))
    box9a.text_frame.paragraphs[0].text = "Machine Learning Pipeline Optimization"
    box9a.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    box9a.name = "Initiative1"

    box9b = slide9.shapes.add_textbox(Inches(1), Inches(3.5), Inches(5), Inches(1.5))
    box9b.text_frame.paragraphs[0].text = "Edge Computing Deployment Strategy"
    box9b.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    box9b.name = "Initiative2"

    # --- Slide 10: Thank You ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[0])
    slide10.shapes.title.text = "Thank You"
    slide10.placeholders[1].text = "Questions? Contact: strategy@acmecorp.com"

    # Save initial pptx (without animations yet)
    prs.save(OUTPUT)
    print(f"Base presentation saved: {OUTPUT}")

    # Now inject animations via XML into slides 3, 5, 7, 9
    inject_animations(OUTPUT)

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


def make_anim_entrance(shape_id, shape_name, delay_ms=0):
    """Create a fly-in entrance animation element for a shape."""
    # This creates a <p:par> element with an appear animation
    anim_xml = f'''<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                         xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <p:cTn id="0" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" nodeType="clickEffect">
            <p:stCondLst>
                <p:cond delay="{delay_ms}"/>
            </p:stCondLst>
            <p:childTnLst>
                <p:set>
                    <p:cBhvr>
                        <p:cTn id="0" dur="1" fill="hold">
                            <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                        </p:cTn>
                        <p:tgtEl>
                            <p:spTgt spid="{shape_id}"/>
                        </p:tgtEl>
                        <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                    </p:cBhvr>
                    <p:to><p:strVal val="visible"/></p:to>
                </p:set>
            </p:childTnLst>
        </p:cTn>
    </p:par>'''
    return ET.fromstring(anim_xml)


def make_anim_emphasis(shape_id, shape_name, delay_ms=0):
    """Create a bold flash emphasis animation element for a shape."""
    anim_xml = f'''<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                         xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <p:cTn id="0" presetID="10" presetClass="emph" presetSubtype="0" fill="hold" nodeType="clickEffect">
            <p:stCondLst>
                <p:cond delay="{delay_ms}"/>
            </p:stCondLst>
            <p:childTnLst>
                <p:animEffect transition="in" filter="fade">
                    <p:cBhvr>
                        <p:cTn id="0" dur="500"/>
                        <p:tgtEl>
                            <p:spTgt spid="{shape_id}"/>
                        </p:tgtEl>
                    </p:cBhvr>
                </p:animEffect>
            </p:childTnLst>
        </p:cTn>
    </p:par>'''
    return ET.fromstring(anim_xml)


def build_timing_element(animations):
    """Build a complete <p:timing> element wrapping animation entries.
    animations: list of ET elements (p:par)
    """
    # Build the sequence that holds all click-triggered animations
    seq_children = []
    for anim in animations:
        seq_children.append(anim)

    timing_xml = '''<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <p:tnLst>
            <p:par>
                <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
                    <p:childTnLst>
                        <p:seq concurrent="1" nextAc="seek">
                            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                                <p:childTnLst>
                                </p:childTnLst>
                            </p:cTn>
                            <p:prevCondLst>
                                <p:cond evt="onPrev" delay="0">
                                    <p:tgtEl><p:sldTgt/></p:tgtEl>
                                </p:cond>
                            </p:prevCondLst>
                            <p:nextCondLst>
                                <p:cond evt="onNext" delay="0">
                                    <p:tgtEl><p:sldTgt/></p:tgtEl>
                                </p:cond>
                            </p:nextCondLst>
                        </p:seq>
                    </p:childTnLst>
                </p:cTn>
            </p:par>
        </p:tnLst>
    </p:timing>'''

    timing = ET.fromstring(timing_xml)
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    child_list = timing.find('.//p:seq/p:cTn/p:childTnLst', ns)
    for anim in animations:
        child_list.append(anim)

    return timing


def get_shape_ids_from_slide_xml(slide_xml_str):
    """Extract shape IDs and names from slide XML."""
    root = ET.fromstring(slide_xml_str)
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
          'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    shapes = []
    # Find all sp (shape) elements with nvSpPr
    for sp in root.iter('{http://schemas.openxmlformats.org/presentationml/2006/main}sp'):
        nvSpPr = sp.find('{http://schemas.openxmlformats.org/presentationml/2006/main}nvSpPr')
        if nvSpPr is not None:
            cNvPr = nvSpPr.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
            if cNvPr is not None:
                sid = cNvPr.get('id')
                name = cNvPr.get('name', '')
                if sid:
                    shapes.append((int(sid), name))
    return shapes


def inject_animations(pptx_path):
    """Post-process the pptx ZIP to inject animation timing XML into slides 3, 5, 7, 9."""
    tmp_path = pptx_path + '.tmp'

    with zipfile.ZipFile(pptx_path, 'r') as zin:
        with zipfile.ZipFile(tmp_path, 'w') as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)

                slide_num = None
                if item.filename.startswith('ppt/slides/slide') and item.filename.endswith('.xml'):
                    try:
                        slide_num = int(item.filename.replace('ppt/slides/slide', '').replace('.xml', ''))
                    except ValueError:
                        pass

                if slide_num in (3, 5, 7, 9):
                    data = add_animations_to_slide(data, slide_num)

                zout.writestr(item, data)

    shutil.move(tmp_path, pptx_path)
    print(f"Animations injected into slides 3, 5, 7, 9")


def add_animations_to_slide(slide_xml_bytes, slide_num):
    """Add animation timing to a slide XML."""
    root = ET.fromstring(slide_xml_bytes)

    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'

    # Get shape IDs from the slide (skip first few that are layout/placeholder)
    shapes = []
    for sp in root.iter(f'{{{ns_p}}}sp'):
        nvSpPr = sp.find(f'{{{ns_p}}}nvSpPr')
        if nvSpPr is not None:
            cNvPr = nvSpPr.find(f'{{{ns_p}}}cNvPr')
            if cNvPr is not None:
                sid = cNvPr.get('id')
                name = cNvPr.get('name', '')
                if sid:
                    shapes.append((int(sid), name))

    if not shapes:
        return slide_xml_bytes

    animations = []

    if slide_num == 3:
        # 2 entrance animations on first 2 shapes
        for i in range(min(2, len(shapes))):
            animations.append(make_anim_entrance(shapes[i][0], shapes[i][1], delay_ms=i*500))

    elif slide_num == 5:
        # 6 animations: 3 entrance + 3 emphasis
        for i in range(min(3, len(shapes))):
            animations.append(make_anim_entrance(shapes[i][0], shapes[i][1], delay_ms=i*500))
        for i in range(3, min(6, len(shapes))):
            animations.append(make_anim_emphasis(shapes[i][0], shapes[i][1], delay_ms=(i-3)*500))

    elif slide_num == 7:
        # 2 entrance animations
        for i in range(min(2, len(shapes))):
            animations.append(make_anim_entrance(shapes[i][0], shapes[i][1], delay_ms=i*500))

    elif slide_num == 9:
        # 2 emphasis animations
        for i in range(min(2, len(shapes))):
            animations.append(make_anim_emphasis(shapes[i][0], shapes[i][1], delay_ms=i*500))

    if animations:
        timing = build_timing_element(animations)
        root.append(timing)

    return ET.tostring(root, xml_declaration=True, encoding='UTF-8')


create_initial()
