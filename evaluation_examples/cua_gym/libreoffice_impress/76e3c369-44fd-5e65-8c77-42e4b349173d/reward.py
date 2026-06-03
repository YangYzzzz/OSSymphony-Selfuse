"""
Reward Script: Remove all animations from slide 5 without affecting other slides
Task ID: impress_gf3_044
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Slide 5 has no timing/animation elements
  Component 2 (0.3): Slide 5 has no animations AND all shapes preserved
  Component 3 (0.3): Slide 5 has no animations AND other slides retain their animations
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_044'

NS = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
}

ANIM_TAGS = ['p:anim', 'p:animEffect', 'p:animMotion', 'p:animRot', 'p:animScale', 'p:set']


def count_animations_in_slide(zf, slide_num):
    """Count animation effect elements in a slide. slide_num is 1-based."""
    try:
        with zf.open(f'ppt/slides/slide{slide_num}.xml') as f:
            root = ET.parse(f).getroot()
            timing = root.find('.//p:timing', NS)
            if timing is None:
                return 0
            count = 0
            for tag in ANIM_TAGS:
                count += len(timing.findall('.//' + tag, NS))
            return count
    except KeyError:
        return -1  # slide doesn't exist


def has_timing_element(zf, slide_num):
    """Check if slide has a timing element with animation content."""
    try:
        with zf.open(f'ppt/slides/slide{slide_num}.xml') as f:
            root = ET.parse(f).getroot()
            timing = root.find('.//p:timing', NS)
            if timing is None:
                return False
            # Check if timing has any meaningful animation children
            for tag in ANIM_TAGS:
                if timing.findall('.//' + tag, NS):
                    return True
            # Also check for seq elements (animation sequences)
            seqs = timing.findall('.//p:seq', NS)
            if seqs:
                return True
            return False
    except KeyError:
        return False


def count_shapes_on_slide(pptx_path, slide_idx):
    """Count shapes on a slide using python-pptx. slide_idx is 0-based."""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        slide = prs.slides[slide_idx]
        return len(slide.shapes)
    except Exception:
        return -1


def get_shape_names(pptx_path, slide_idx):
    """Get shape names on a slide. slide_idx is 0-based."""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        slide = prs.slides[slide_idx]
        return [shape.name for shape in slide.shapes]
    except Exception:
        return []


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: file must be loadable
    if not os.path.exists(file_path):
        print(f"CRITICAL: File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    try:
        zf = zipfile.ZipFile(file_path, 'r')
    except Exception as e:
        print(f"CRITICAL: Cannot open zip: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have 10 slides
    slide_files = [f for f in zf.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
    if len(slide_files) != 10:
        print(f"CRITICAL: Expected 10 slides, found {len(slide_files)}")
        zf.close()
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide 5 has NO animation effects (0.4 points)
    # Initial has 6 animation effects on slide 5; golden should have 0.
    try:
        slide5_anims = count_animations_in_slide(zf, 5)
        slide5_has_timing = has_timing_element(zf, 5)
        if slide5_anims == 0 and not slide5_has_timing:
            print(f"PASS: Component 1 — Slide 5 has no animations and no timing element (0.4 pts)")
            total_score += 0.4
        elif slide5_anims == 0:
            # Timing element exists but no actual animations — partial credit
            print(f"PASS (partial): Component 1 — Slide 5 has no animation effects but timing element remains (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 — Slide 5 still has {slide5_anims} animation effects")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 5 has no animations AND all 8 shapes are preserved (0.3 points)
    # This is a compound check: animation removal is the task change, shape preservation is correctness.
    # Anchored to animation change so it fails on initial_env.
    try:
        slide5_anims_c2 = count_animations_in_slide(zf, 5)
        shape_names = get_shape_names(file_path, 4)  # 0-indexed
        expected_names = ['Title 1', 'RoadmapTitle', 'Phase1', 'Phase2', 'Phase3', 'Budget', 'ROI', 'Metric']

        anims_removed = (slide5_anims_c2 == 0)
        shapes_ok = (len(shape_names) >= 8)
        names_match = all(name in shape_names for name in expected_names)

        if anims_removed and shapes_ok and names_match:
            print(f"PASS: Component 2 — Slide 5 animations removed AND all 8 shapes preserved (0.3 pts)")
            total_score += 0.3
        elif anims_removed and shapes_ok:
            print(f"PASS (partial): Component 2 — Animations removed, {len(shape_names)} shapes found but some names differ (0.2 pts)")
            total_score += 0.2
        elif anims_removed:
            print(f"FAIL: Component 2 — Animations removed but only {len(shape_names)} shapes remain (expected 8)")
        else:
            print(f"FAIL: Component 2 — Slide 5 still has {slide5_anims_c2} animations")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 5 has no animations AND other slides (3, 7, 9) retain their animations (0.3 points)
    # This compound check ensures animations were only removed from slide 5.
    # Anchored to animation change on slide 5 so it fails on initial_env.
    try:
        slide5_anims_c3 = count_animations_in_slide(zf, 5)
        slide3_anims = count_animations_in_slide(zf, 3)
        slide7_anims = count_animations_in_slide(zf, 7)
        slide9_anims = count_animations_in_slide(zf, 9)

        slide5_clear = (slide5_anims_c3 == 0)
        others_intact = (slide3_anims >= 2 and slide7_anims >= 2 and slide9_anims >= 2)

        if slide5_clear and others_intact:
            print(f"PASS: Component 3 — Slide 5 clear AND slides 3({slide3_anims}), 7({slide7_anims}), 9({slide9_anims}) retain animations (0.3 pts)")
            total_score += 0.3
        elif slide5_clear and (slide3_anims >= 2 or slide7_anims >= 2 or slide9_anims >= 2):
            intact_count = sum(1 for a in [slide3_anims, slide7_anims, slide9_anims] if a >= 2)
            partial = round(0.1 * intact_count, 1)
            print(f"PARTIAL: Component 3 — Slide 5 clear, but only {intact_count}/3 other slides retain animations ({partial} pts)")
            total_score += partial
        elif not slide5_clear:
            print(f"FAIL: Component 3 — Slide 5 still has {slide5_anims_c3} animations")
        else:
            print(f"FAIL: Component 3 — Other slides lost animations: s3={slide3_anims}, s7={slide7_anims}, s9={slide9_anims}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    zf.close()

    final_score = min(round(total_score, 1), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
