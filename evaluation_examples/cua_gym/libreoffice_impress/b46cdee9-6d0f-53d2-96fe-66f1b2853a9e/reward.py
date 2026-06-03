"""
Reward Script: Add grouped bar chart to slide 6 of Annual_Review_2025.pptx
Task ID: impress_ps_020
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Chart exists on slide 6 and is a bar/column clustered type
  Component 2 (0.25): Correct 4 categories and 3 series
  Component 3 (0.30): Correct data values in all series
  Component 4 (0.10): Legend is present
  Component 5 (0.10): Axes visible and series have distinct fill colors
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_020'


def persist_app_state(domain):
    """Try to save any unsaved LibreOffice state via Ctrl+S."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation
    from pptx.enum.chart import XL_CHART_TYPE

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[5]  # 0-indexed, slide 6

    # Find chart shape on slide 6
    chart_shape = None
    for shape in slide.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Chart exists on slide 6 and is bar/column clustered (0.25 points)
    try:
        if chart_shape is None:
            print("FAIL: Component 1 — No chart found on slide 6")
        else:
            chart = chart_shape.chart
            chart_type = chart.chart_type
            # Accept BAR_CLUSTERED (57) or COLUMN_CLUSTERED (51) as valid grouped bar charts
            valid_types = [XL_CHART_TYPE.BAR_CLUSTERED, XL_CHART_TYPE.COLUMN_CLUSTERED]
            if chart_type in valid_types:
                print(f"PASS: Component 1 — Grouped bar chart found on slide 6, type={chart_type} (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 — Chart on slide 6 is type {chart_type}, expected BAR_CLUSTERED or COLUMN_CLUSTERED")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    if chart_shape is None:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    chart = chart_shape.chart
    plot = chart.plots[0]

    # Component 2: Correct categories (4 departments) and 3 series (0.25 points)
    try:
        expected_categories = ['Engineering', 'Sales', 'Marketing', 'Operations']
        try:
            actual_categories = [str(c) for c in plot.categories]
        except Exception:
            actual_categories = []

        num_series = len(plot.series)
        cat_match = (actual_categories == expected_categories)
        series_match = (num_series == 3)

        if cat_match and series_match:
            print(f"PASS: Component 2 — 4 correct categories and 3 series (0.25 pts)")
            total_score += 0.25
        else:
            if not cat_match:
                print(f"FAIL: Component 2 — Categories: expected {expected_categories}, found {actual_categories}")
            if not series_match:
                print(f"FAIL: Component 2 — Series count: expected 3, found {num_series}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Correct data values in all 3 series (0.30 points)
    # Expected: Engineering (95, 4.2, 88), Sales (110, 4.5, 92), Marketing (85, 3.9, 78), Operations (92, 4.0, 95)
    # Series 0 (Target Achievement %): [95, 110, 85, 92]
    # Series 1 (Customer Satisfaction): [4.2, 4.5, 3.9, 4.0]
    # Series 2 (Budget Utilization %): [88, 92, 78, 95]
    try:
        expected_data = [
            [95.0, 110.0, 85.0, 92.0],
            [4.2, 4.5, 3.9, 4.0],
            [88.0, 92.0, 78.0, 95.0],
        ]

        series_correct = 0
        for i, expected_vals in enumerate(expected_data):
            if i < len(plot.series):
                actual_vals = list(plot.series[i].values)
                if len(actual_vals) == len(expected_vals):
                    all_match = all(
                        abs(a - e) < 0.01
                        for a, e in zip(actual_vals, expected_vals)
                    )
                    if all_match:
                        series_correct += 1
                        print(f"  Series {i} data: CORRECT {actual_vals}")
                    else:
                        print(f"  Series {i} data: MISMATCH expected {expected_vals}, got {actual_vals}")
                else:
                    print(f"  Series {i} data: LENGTH MISMATCH expected {len(expected_vals)}, got {len(actual_vals)}")
            else:
                print(f"  Series {i}: MISSING")

        # Award proportional credit: 0.10 per correct series
        points = round(series_correct * 0.10, 2)
        if series_correct == 3:
            print(f"PASS: Component 3 — All 3 series have correct data values (0.30 pts)")
            total_score += 0.30
        elif series_correct > 0:
            print(f"PARTIAL: Component 3 — {series_correct}/3 series correct ({points} pts)")
            total_score += points
        else:
            print("FAIL: Component 3 — No series data matches expected values")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Legend is present (0.10 points)
    try:
        if chart.has_legend:
            print("PASS: Component 4 — Chart has legend (0.10 pts)")
            total_score += 0.10
        else:
            print("FAIL: Component 4 — Chart has no legend")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Axes visible and series have distinct colors (0.10 points)
    try:
        axes_ok = False
        colors_ok = False

        # Check axes visibility
        try:
            cat_visible = chart.category_axis.visible
            val_visible = chart.value_axis.visible
            axes_ok = (cat_visible and val_visible)
            if axes_ok:
                print("  Axes: both category and value axis visible")
            else:
                print(f"  Axes: category={cat_visible}, value={val_visible}")
        except Exception as e:
            print(f"  Axes check error: {e}")

        # Check distinct colors across series
        try:
            colors = set()
            for i, series in enumerate(plot.series):
                fill = series.format.fill
                if fill.type is not None:
                    try:
                        rgb = str(fill.fore_color.rgb)
                        colors.add(rgb)
                    except Exception:
                        pass
            colors_ok = (len(colors) >= 3)
            if colors_ok:
                print(f"  Colors: {len(colors)} distinct series colors found: {colors}")
            else:
                print(f"  Colors: only {len(colors)} distinct colors found, need >= 3")
        except Exception as e:
            print(f"  Colors check error: {e}")

        if axes_ok and colors_ok:
            print("PASS: Component 5 — Axes visible and distinct colors (0.10 pts)")
            total_score += 0.10
        elif axes_ok or colors_ok:
            print(f"PARTIAL: Component 5 — axes_ok={axes_ok}, colors_ok={colors_ok} (0.05 pts)")
            total_score += 0.05
        else:
            print("FAIL: Component 5 — Neither axes visible nor distinct colors")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
