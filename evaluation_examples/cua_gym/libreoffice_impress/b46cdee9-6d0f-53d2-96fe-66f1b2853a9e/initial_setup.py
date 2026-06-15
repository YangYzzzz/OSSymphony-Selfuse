"""
Initial Setup: Create Annual_Review_2025.pptx with 12 slides.
Slide 6 has title 'Department Performance' and empty content area. No charts.
Task ID: impress_ps_020
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_020'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.paragraphs[0].text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
    return slide


def add_title_only_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(prs, "Annual Performance Review 2025",
                    "Prepared by Strategic Planning Division\nConfidential — Internal Use Only")

    # Slide 2: Agenda
    add_content_slide(prs, "Review Agenda", [
        "Executive Summary & Key Highlights",
        "Financial Performance Overview",
        "Revenue Breakdown by Region",
        "Workforce & Talent Analytics",
        "Department Performance Comparison",
        "Strategic Initiatives & Milestones",
        "Customer Engagement Metrics",
        "Technology & Infrastructure Updates",
        "Risk Assessment & Mitigation",
        "2026 Goals & Roadmap",
    ])

    # Slide 3: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "Total revenue grew 18.3% year-over-year to $247.5M",
        "Net profit margin improved from 12.1% to 14.7%",
        "Employee satisfaction score reached 4.3/5.0 (up from 3.8)",
        "Customer retention rate at 94.2%, exceeding 90% target",
        "Successfully launched 3 new product lines in Q2 and Q3",
        "Expanded operations to 4 new international markets",
    ])

    # Slide 4: Financial Overview
    add_content_slide(prs, "Financial Performance Overview", [
        "Q1 Revenue: $54.2M | Q2: $61.8M | Q3: $67.1M | Q4: $64.4M",
        "Operating expenses reduced by 7.2% through process optimization",
        "Capital expenditure: $18.9M invested in infrastructure upgrades",
        "Cash reserves increased to $92.3M, providing 14 months runway",
        "Debt-to-equity ratio improved from 0.45 to 0.32",
    ])

    # Slide 5: Revenue by Region
    add_content_slide(prs, "Revenue Breakdown by Region", [
        "North America: $142.8M (57.7%) — up 15.4%",
        "Europe & UK: $58.3M (23.6%) — up 22.1%",
        "Asia-Pacific: $31.2M (12.6%) — up 28.9%",
        "Latin America: $9.8M (4.0%) — up 11.3%",
        "Middle East & Africa: $5.4M (2.1%) — new market entry",
    ])

    # Slide 6: Department Performance — TITLE ONLY, EMPTY CONTENT AREA
    slide6 = add_title_only_slide(prs, "Department Performance")
    # Intentionally empty: no chart, no data, no content placeholder beyond title

    # Slide 7: Strategic Initiatives
    add_content_slide(prs, "Strategic Initiatives & Milestones", [
        "Project Aurora: Cloud migration completed ahead of schedule (Q2)",
        "Project Horizon: New CRM platform rolled out to 1,200 users",
        "Project Catalyst: AI-powered analytics dashboard launched in Q3",
        "Sustainability program reduced carbon footprint by 23%",
        "ISO 27001 certification achieved in September 2025",
    ])

    # Slide 8: Customer Engagement
    add_content_slide(prs, "Customer Engagement Metrics", [
        "Net Promoter Score: 72 (industry average: 48)",
        "Average response time reduced from 4.2 hours to 1.8 hours",
        "Support ticket resolution rate: 96.7% within SLA",
        "Monthly active users grew 34% to 1.2M",
        "Customer onboarding time reduced by 40%",
    ])

    # Slide 9: Workforce Analytics
    add_content_slide(prs, "Workforce & Talent Analytics", [
        "Total headcount: 2,847 (net increase of 312)",
        "Voluntary turnover rate: 8.3% (industry average: 13.2%)",
        "Diversity hiring increased 28% across all departments",
        "Training investment: $4.2M — 42 hours average per employee",
        "Internal promotion rate: 34% of all new roles filled internally",
    ])

    # Slide 10: Technology Updates
    add_content_slide(prs, "Technology & Infrastructure Updates", [
        "System uptime: 99.97% across all production environments",
        "Migrated 78% of workloads to multi-cloud architecture",
        "Cybersecurity: Zero critical breaches; 12 threats neutralized",
        "DevOps pipeline reduced deployment time from 4 hours to 22 minutes",
        "Data warehouse capacity expanded by 340% for analytics growth",
    ])

    # Slide 11: Risk Assessment
    add_content_slide(prs, "Risk Assessment & Mitigation", [
        "Supply chain diversification: reduced single-vendor dependency to 15%",
        "Regulatory compliance: passed all 7 audits with zero findings",
        "Business continuity: DR failover tested quarterly, RTO < 2 hours",
        "Currency hedging strategy saved $3.1M in FX exposure",
        "Talent pipeline: 6-month succession plans for all VP+ roles",
    ])

    # Slide 12: 2026 Goals
    add_content_slide(prs, "2026 Goals & Strategic Roadmap", [
        "Revenue target: $295M (19.2% growth)",
        "Expand to 3 additional international markets",
        "Launch next-gen product platform (Project Stellar) in Q2",
        "Achieve carbon-neutral operations by Q4 2026",
        "Increase employee NPS from 62 to 75",
        "Invest $25M in R&D for emerging technologies",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
