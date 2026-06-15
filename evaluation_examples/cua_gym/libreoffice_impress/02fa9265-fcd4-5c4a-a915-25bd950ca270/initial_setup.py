"""
Initial Setup: Add slide number counter to master slide
Task ID: impress_gf3_037
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_037'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # Set standard 16:9 slide size
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Set dark background on slide master ---
    master = prs.slide_masters[0]
    fill = master.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)  # Dark navy background

    # Make default text white on master layouts
    # We'll set text colors per-slide since master text styling is complex

    # Slide topics for a realistic 18-slide deck about "Q4 2025 Strategic Planning"
    slide_content = [
        ("Q4 2025 Strategic Planning", "Executive Leadership Summit\nNovember 2025"),
        ("Agenda", "• Financial Overview\n• Market Analysis\n• Product Roadmap\n• Talent Strategy\n• Risk Assessment"),
        ("Financial Overview", "Revenue: $142.3M (+18% YoY)\nEBITDA: $28.7M (20.2% margin)\nFree Cash Flow: $19.4M"),
        ("Revenue Breakdown by Segment", "Enterprise: $78.5M (55%)\nMid-Market: $42.1M (30%)\nSMB: $21.7M (15%)"),
        ("Cost Structure Analysis", "COGS: $68.3M (48%)\nR&D: $21.3M (15%)\nS&M: $28.5M (20%)\nG&A: $14.2M (10%)"),
        ("Market Analysis", "TAM Growth: 12.4% CAGR\nOur Market Share: 8.7% → 11.2%\nKey Competitor Moves: 3 new entrants"),
        ("Customer Metrics", "NPS Score: 72 (+8)\nChurn Rate: 4.2% (-1.1%)\nLTV/CAC Ratio: 4.8x"),
        ("Product Roadmap", "Phase 1: AI Integration (Dec 2025)\nPhase 2: Platform Expansion (Feb 2026)\nPhase 3: Enterprise SSO (Apr 2026)"),
        ("Engineering Velocity", "Sprint Velocity: 142 pts/sprint\nDeployment Frequency: 4.2/week\nIncident Rate: 0.3% (-40% QoQ)"),
        ("Talent Strategy", "Current Headcount: 487\nOpen Roles: 34\nAttrition Rate: 8.2%\nDEI Index: 0.74"),
        ("Hiring Pipeline", "Engineering: 18 roles\nProduct: 6 roles\nSales: 7 roles\nOperations: 3 roles"),
        ("Compensation Benchmarking", "Engineering: P65 → P72\nProduct: P58 → P65\nSales OTE: +12% adjustment\nEquity refresh: 15% pool increase"),
        ("Risk Assessment", "Regulatory: Medium (GDPR expansion)\nTechnical: Low (redundancy in place)\nMarket: Medium (price pressure)\nTalent: High (competitive market)"),
        ("Mitigation Strategies", "• Legal team expansion for compliance\n• Multi-cloud infrastructure\n• Value-based pricing model\n• Retention bonuses for key talent"),
        ("Q4 OKRs", "O1: Achieve $38M quarterly revenue\nO2: Launch AI features to 50% of users\nO3: Reduce churn below 3.5%\nO4: Fill all critical roles by Dec 15"),
        ("Budget Allocation", "Engineering: $5.8M\nMarketing: $3.2M\nSales: $2.9M\nInfrastructure: $1.7M\nTotal: $13.6M"),
        ("Timeline & Milestones", "Oct 1: Q4 kickoff\nOct 31: AI beta launch\nNov 15: Mid-quarter review\nDec 15: Hiring complete\nDec 31: Q4 close"),
        ("Next Steps & Actions", "1. Department leads submit detailed plans by Oct 7\n2. Budget finalization meeting Oct 10\n3. All-hands Q4 kickoff Oct 15\n4. Weekly exec sync every Monday"),
    ]

    for i, (title_text, body_text) in enumerate(slide_content):
        if i == 0:
            # Title slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = title_text
            for run in title_shape.text_frame.paragraphs[0].runs:
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.size = Pt(36)
                run.font.bold = True

            # Set subtitle
            subtitle = slide.placeholders[1]
            subtitle.text = body_text
            for para in subtitle.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
                    run.font.size = Pt(18)
        else:
            # Content slides
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            # Title
            title_shape = slide.shapes.title
            title_shape.text = title_text
            for run in title_shape.text_frame.paragraphs[0].runs:
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.size = Pt(28)
                run.font.bold = True

            # Body content
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            lines = body_text.split('\n')
            for j, line in enumerate(lines):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                for run in p.runs:
                    run.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
                    run.font.size = Pt(16)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
