"""
Initial Setup: Create interactive quiz presentation with 10 slides
Task ID: impress_gf2_004
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_004'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_centered_text(slide, text, left, top, width, height,
                      font_size=28, bold=False, color=None, alignment=PP_ALIGN.CENTER):
    """Add a textbox with centered text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_button_shape(slide, text, left, top, width, height,
                     fill_color, font_size=20, font_color=RGBColor(0xFF, 0xFF, 0xFF)):
    """Add a rounded rectangle button shape with text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # no border

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = font_color
    return shape


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== Slide 1: Title Slide =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    add_centered_text(slide1, "Geography Quiz Challenge", Inches(2), Inches(1.5),
                      Inches(9), Inches(1.5), font_size=44, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF))
    add_centered_text(slide1, "Test your knowledge of world capitals!",
                      Inches(3), Inches(3.5), Inches(7), Inches(1),
                      font_size=24, color=RGBColor(0xCC, 0xDD, 0xEE))
    add_centered_text(slide1, "Click to Begin", Inches(4), Inches(5.5),
                      Inches(5), Inches(1), font_size=20,
                      color=RGBColor(0x88, 0xBB, 0xDD))

    # ===== Slide 2: Instructions =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = slide2.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    add_centered_text(slide2, "How to Play", Inches(3), Inches(0.8),
                      Inches(7), Inches(1), font_size=36, bold=True,
                      color=RGBColor(0x1B, 0x3A, 0x5C))
    instructions = [
        "1. Read each question carefully",
        "2. Click on the answer button you think is correct",
        "3. You will be told if your answer is right or wrong",
        "4. Try to get all questions correct!"
    ]
    y_pos = Inches(2.5)
    for inst in instructions:
        add_centered_text(slide2, inst, Inches(2), y_pos, Inches(9), Inches(0.6),
                          font_size=20, color=RGBColor(0x33, 0x33, 0x33),
                          alignment=PP_ALIGN.LEFT)
        y_pos += Inches(0.8)

    # ===== Slide 3: Category Intro =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = slide3.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    add_centered_text(slide3, "Round 1: European Capitals",
                      Inches(2), Inches(2.5), Inches(9), Inches(1.5),
                      font_size=40, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_centered_text(slide3, "Get ready...", Inches(4), Inches(4.5),
                      Inches(5), Inches(1), font_size=22,
                      color=RGBColor(0xAA, 0xCC, 0xDD))

    # ===== Slide 4: Quiz Question (THE KEY SLIDE) =====
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = slide4.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

    # Question text
    add_centered_text(slide4, "What is the capital of France?",
                      Inches(2), Inches(0.8), Inches(9), Inches(1.2),
                      font_size=36, bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))

    # Three answer buttons (rectangle shapes, no hyperlinks)
    btn_width = Inches(8)
    btn_height = Inches(1.0)
    btn_left = Inches(2.6)

    add_button_shape(slide4, "A: London", btn_left, Inches(2.8),
                     btn_width, btn_height,
                     fill_color=RGBColor(0x34, 0x98, 0xDB))

    add_button_shape(slide4, "B: Paris", btn_left, Inches(4.2),
                     btn_width, btn_height,
                     fill_color=RGBColor(0x27, 0xAE, 0x60))

    add_button_shape(slide4, "C: Berlin", btn_left, Inches(5.6),
                     btn_width, btn_height,
                     fill_color=RGBColor(0xE7, 0x4C, 0x3C))

    # ===== Slide 5: Filler - Quiz Facts =====
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = slide5.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xEC, 0xF0, 0xF1)
    add_centered_text(slide5, "Did You Know?", Inches(3), Inches(1),
                      Inches(7), Inches(1), font_size=32, bold=True,
                      color=RGBColor(0x2C, 0x3E, 0x50))
    add_centered_text(slide5, "Paris has been the capital of France since 987 AD,\n"
                      "making it one of the oldest capital cities in Europe.",
                      Inches(2), Inches(3), Inches(9), Inches(2),
                      font_size=20, color=RGBColor(0x55, 0x55, 0x55))

    # ===== Slide 6: Correct Answer =====
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = slide6.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x27, 0xAE, 0x60)
    add_centered_text(slide6, "Correct!", Inches(3), Inches(2.5),
                      Inches(7), Inches(1.5), font_size=48, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF))
    # No "Next Question" button in initial state

    # ===== Slide 7: Wrong Answer =====
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = slide7.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xE7, 0x4C, 0x3C)
    add_centered_text(slide7, "Try Again!", Inches(3), Inches(2.5),
                      Inches(7), Inches(1.5), font_size=48, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF))

    # ===== Slide 8: Next Question Placeholder =====
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = slide8.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    add_centered_text(slide8, "Question 2: Coming Soon!",
                      Inches(2), Inches(2.5), Inches(9), Inches(1.5),
                      font_size=36, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_centered_text(slide8, "What is the largest ocean on Earth?",
                      Inches(3), Inches(4.5), Inches(7), Inches(1),
                      font_size=24, color=RGBColor(0xCC, 0xDD, 0xEE))

    # ===== Slide 9: Score Tracking =====
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = slide9.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    add_centered_text(slide9, "Score Tracker", Inches(3), Inches(1),
                      Inches(7), Inches(1), font_size=32, bold=True,
                      color=RGBColor(0x1B, 0x3A, 0x5C))
    add_centered_text(slide9, "Keep track of your correct answers\nas you progress through the quiz.",
                      Inches(2), Inches(3), Inches(9), Inches(2),
                      font_size=20, color=RGBColor(0x55, 0x55, 0x55))

    # ===== Slide 10: End Slide =====
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = slide10.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    add_centered_text(slide10, "Quiz Complete!", Inches(3), Inches(2),
                      Inches(7), Inches(1.5), font_size=44, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF))
    add_centered_text(slide10, "Thank you for playing the Geography Quiz Challenge!",
                      Inches(2), Inches(4), Inches(9), Inches(1),
                      font_size=22, color=RGBColor(0xCC, 0xDD, 0xEE))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
