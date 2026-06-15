"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 104 I need Picture 1 to sit exactly 2.0 cm from the left edge and 6.0 cm down from the top. What’s the quickest way to set those precise coordinates in LibreOffice Impress?
Generated: 2025-09-10 16:04:33
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
import glob
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def verify_picture_coordinates(file_path: str,
                               slide_index: int = 103,
                               expected_x_cm: float = 2.0,
                               expected_y_cm: float = 6.0,
                               tolerance_cm: float = 0.05) -> float:
    """Verify that on the specified slide the *first* picture (Picture 1)
    sits exactly at the expected coordinates within a tolerance.

    A progressive score (0.0 – 1.0) is returned based on:
      • File successfully loaded   (no points – prerequisite)
      • Slide 104 exists           (0.2)
      • A picture exists           (0.3)
      • Correct X coordinate       (0.25)
      • Correct Y coordinate       (0.25)
    Total possible = 1.0
    """
    score = 0.0
    max_score = 1.0

    print(f"Loading presentation: {file_path}")

    # ---------- 1. Load presentation (prerequisite – no points) ----------
    if not os.path.exists(file_path):
        print("✗ File not found – cannot verify task")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error loading PPTX: {e}")
        return 0.0

    # ---------- 2. Verify slide 104 exists (index 103) ----------
    if len(prs.slides) > slide_index:
        print(f"✓ Slide count is {len(prs.slides)} (>= 104)")
        score += 0.2
    else:
        print(f"✗ Slide 104 missing (only {len(prs.slides)} slides)")
        return score  # cannot proceed further but return partial progress

    slide = prs.slides[slide_index]

    # ---------- 3. Locate first picture on the slide ----------
    pictures = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if not pictures:
        print("✗ No picture shapes found on slide 104")
        return score

    pic = pictures[0]
    print(f"✓ Found picture shape: name=\"{pic.name}\"")
    score += 0.3

    # ---------- 4. Measure coordinates ----------
    EMU_PER_CM = 914400 / 2.54  # 1 inch = 914400 EMU; 2.54 cm per inch
    x_cm = pic.left / EMU_PER_CM
    y_cm = pic.top / EMU_PER_CM

    print(f"Picture position: {x_cm:.2f} cm from left, {y_cm:.2f} cm from top")

    # ---------- 5. Evaluate X coordinate ----------
    if abs(x_cm - expected_x_cm) <= tolerance_cm:
        print("✓ X coordinate within tolerance")
        score += 0.25
    else:
        print("✗ X coordinate incorrect")

    # ---------- 6. Evaluate Y coordinate ----------
    if abs(y_cm - expected_y_cm) <= tolerance_cm:
        print("✓ Y coordinate within tolerance")
        score += 0.25
    else:
        print("✗ Y coordinate incorrect")

    final_score = min(score, max_score)
    print(f"REWARD: {final_score}")
    return final_score


# -------------------- Script Execution --------------------
if __name__ == "__main__":
    # Attempt to locate the relevant .pptx in the user's home directory
    target_path = None
    for p in glob.glob(os.path.expanduser("~/" + "*.pptx")):
        if "on_slide_104" in os.path.basename(p):
            target_path = p
            break
    # Fallback to any pptx if the specific naming pattern is not found
    if target_path is None:
        matches = glob.glob(os.path.expanduser("~/" + "*.pptx"))
        if matches:
            target_path = matches[0]

    if target_path:
        verify_picture_coordinates(target_path)
    else:
        print("✗ No .pptx file found for verification")
