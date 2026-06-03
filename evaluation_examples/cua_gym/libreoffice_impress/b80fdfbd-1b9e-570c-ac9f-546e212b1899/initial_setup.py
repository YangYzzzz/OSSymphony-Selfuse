"""
Initial Setup: Build a presentation with 8 slides, slide 3 titled 'By The Numbers' with empty content.
Task ID: impress_sales_086
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_086'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_name="Arial",
                font_size=Pt(18), bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a styled text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== Slide 1: Title Slide =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only
    add_textbox(slide1, Inches(1.5), Inches(1.5), Inches(10), Inches(1.5),
                "Numbers Pitch Deck", "Arial", Pt(44), True,
                RGBColor(0x1A, 0x3C, 0x6E), PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(2), Inches(3.5), Inches(9), Inches(1),
                "Q4 2025 Performance Review — NovaTech Solutions", "Arial", Pt(22), False,
                RGBColor(0x55, 0x55, 0x55), PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(3), Inches(5.0), Inches(7), Inches(0.6),
                "Presented by Sarah Chen, VP of Product", "Arial", Pt(16), False,
                RGBColor(0x88, 0x88, 0x88), PP_ALIGN.CENTER)

    # ===== Slide 2: Our Mission =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide2, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                "Our Mission", "Arial", Pt(36), True,
                RGBColor(0x1A, 0x3C, 0x6E), PP_ALIGN.LEFT)
    mission_text = (
        "NovaTech Solutions empowers businesses of all sizes to streamline their "
        "operations through intelligent automation and data-driven insights. Since "
        "our founding in 2019, we have helped over 2,500 companies reduce operational "
        "costs by an average of 32% while improving customer satisfaction scores."
    )
    add_textbox(slide2, Inches(0.8), Inches(1.8), Inches(11), Inches(2),
                mission_text, "Arial", Pt(18), False,
                RGBColor(0x33, 0x33, 0x33), PP_ALIGN.LEFT)
    add_textbox(slide2, Inches(0.8), Inches(4.2), Inches(11), Inches(1.5),
                "Our platform integrates seamlessly with existing enterprise tools, "
                "providing real-time analytics dashboards, automated workflow management, "
                "and predictive maintenance alerts that keep your business running at peak efficiency.",
                "Arial", Pt(18), False, RGBColor(0x33, 0x33, 0x33), PP_ALIGN.LEFT)

    # ===== Slide 3: By The Numbers (EMPTY content — task target) =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide3, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                "By The Numbers", "Arial", Pt(36), True,
                RGBColor(0x1A, 0x3C, 0x6E), PP_ALIGN.LEFT)
    # Content area intentionally left empty — the agent must add stats here

    # ===== Slide 4: Revenue Breakdown =====
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide4, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                "Revenue Breakdown", "Arial", Pt(36), True,
                RGBColor(0x1A, 0x3C, 0x6E), PP_ALIGN.LEFT)
    revenue_items = [
        ("Enterprise Licenses", "$28.5M", "57% of total revenue"),
        ("SMB Subscriptions", "$12.3M", "25% of total revenue"),
        ("Professional Services", "$6.2M", "12% of total revenue"),
        ("Training & Support", "$3.0M", "6% of total revenue"),
    ]
    y_pos = 1.8
    for label, amount, pct in revenue_items:
        add_textbox(slide4, Inches(1.2), Inches(y_pos), Inches(4), Inches(0.5),
                    f"{label}: {amount}", "Arial", Pt(20), True,
                    RGBColor(0x2B, 0x6C, 0xB0), PP_ALIGN.LEFT)
        add_textbox(slide4, Inches(6), Inches(y_pos), Inches(5), Inches(0.5),
                    pct, "Arial", Pt(16), False,
                    RGBColor(0x66, 0x66, 0x66), PP_ALIGN.LEFT)
        y_pos += 1.1

    # ===== Slide 5: Customer Growth =====
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide5, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                "Customer Growth Trajectory", "Arial", Pt(36), True,
                RGBColor(0x1A, 0x3C, 0x6E), PP_ALIGN.LEFT)
    growth_data = [
        ("Q1 2025", "8,200 active users", "+12% QoQ"),
        ("Q2 2025", "8,950 active users", "+9% QoQ"),
        ("Q3 2025", "9,600 active users", "+7% QoQ"),
        ("Q4 2025", "10,340 active users", "+8% QoQ"),
    ]
    y_pos = 1.8
    for quarter, users, growth in growth_data:
        add_textbox(slide5, Inches(1.2), Inches(y_pos), Inches(3), Inches(0.5),
                    quarter, "Arial", Pt(20), True,
                    RGBColor(0x33, 0x33, 0x33), PP_ALIGN.LEFT)
        add_textbox(slide5, Inches(4.5), Inches(y_pos), Inches(4), Inches(0.5),
                    users, "Arial", Pt(18), False,
                    RGBColor(0x2B, 0x6C, 0xB0), PP_ALIGN.LEFT)
        add_textbox(slide5, Inches(9), Inches(y_pos), Inches(3), Inches(0.5),
                    growth, "Arial", Pt(16), False,
                    RGBColor(0x27, 0xAE, 0x60), PP_ALIGN.LEFT)
        y_pos += 1.1

    # ===== Slide 6: Product Roadmap =====
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide6, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                "2026 Product Roadmap", "Arial", Pt(36), True,
                RGBColor(0x1A, 0x3C, 0x6E), PP_ALIGN.LEFT)
    roadmap_items = [
        "Q1: AI-Powered Analytics Dashboard v3.0 with natural language queries",
        "Q2: Mobile Application Launch (iOS & Android) with offline capability",
        "Q3: Enterprise SSO Integration & Advanced Role-Based Access Control",
        "Q4: Predictive Revenue Forecasting Module with ML-driven insights",
    ]
    y_pos = 1.8
    for item in roadmap_items:
        add_textbox(slide6, Inches(1.2), Inches(y_pos), Inches(10), Inches(0.7),
                    item, "Arial", Pt(18), False,
                    RGBColor(0x33, 0x33, 0x33), PP_ALIGN.LEFT)
        y_pos += 1.1

    # ===== Slide 7: Team Highlights =====
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide7, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                "Team Highlights", "Arial", Pt(36), True,
                RGBColor(0x1A, 0x3C, 0x6E), PP_ALIGN.LEFT)
    team_info = [
        "Engineering: 45 engineers across 3 offices (SF, Austin, Berlin)",
        "Customer Success: 98% renewal rate, NPS score of 72",
        "Sales: 15 enterprise deals closed in Q4, avg deal size $185K",
        "Support: Average response time under 4 hours, 24/7 coverage",
    ]
    y_pos = 1.8
    for info in team_info:
        add_textbox(slide7, Inches(1.2), Inches(y_pos), Inches(10), Inches(0.7),
                    info, "Arial", Pt(18), False,
                    RGBColor(0x33, 0x33, 0x33), PP_ALIGN.LEFT)
        y_pos += 1.1

    # ===== Slide 8: Thank You =====
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide8, Inches(2), Inches(2.0), Inches(9), Inches(1.5),
                "Thank You", "Arial", Pt(48), True,
                RGBColor(0x1A, 0x3C, 0x6E), PP_ALIGN.CENTER)
    add_textbox(slide8, Inches(2), Inches(4.0), Inches(9), Inches(1),
                "Questions? Contact us at investor-relations@novatech.io", "Arial", Pt(20), False,
                RGBColor(0x55, 0x55, 0x55), PP_ALIGN.CENTER)
    add_textbox(slide8, Inches(2), Inches(5.5), Inches(9), Inches(0.6),
                "NovaTech Solutions  |  www.novatech.io  |  (415) 555-0192", "Arial", Pt(14), False,
                RGBColor(0x88, 0x88, 0x88), PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
