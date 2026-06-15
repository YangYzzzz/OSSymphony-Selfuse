"""
Reward Script: Build 'By The Numbers' slide with four animated statistics
Task ID: impress_sales_086
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Four stat text boxes with correct text on slide 3
  Component 2 (0.20): Four subtitle text boxes with correct text on slide 3
  Component 3 (0.20): Stat formatting: 48pt bold #2B6CB0
  Component 4 (0.15): Subtitle formatting: 16pt #666666
  Component 5 (0.15): Zoom entrance animations (presetID=53) present on slide 3
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_086'

# Expected stats and subtitles
EXPECTED_STATS = ['10,000+', '$50M+', '99.9%', '4.8/5']
EXPECTED_SUBTITLES = ['Active Users', 'Revenue Managed', 'Uptime', 'Customer Rating']


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]  # 0-indexed

    # Gather all text shapes on slide 3 (excluding title)
    shapes_info = []
    for shape in slide3.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    # Get font properties from runs
                    runs = [r for r in para.runs if (r.text or "").strip()]
                    if runs:
                        run = runs[0]
                        font_size = run.font.size  # EMU
                        font_bold = run.font.bold
                        try:
                            font_color = str(run.font.color.rgb) if run.font.color.type is not None else None
                        except:
                            font_color = None
                        shapes_info.append({
                            'text': text,
                            'size': font_size,
                            'bold': font_bold,
                            'color': font_color,
                        })

    print(f"INFO: Found {len(shapes_info)} text items on slide 3")
    for si in shapes_info:
        print(f"  - text='{si['text']}' size={si['size']} bold={si['bold']} color={si['color']}")

    # Component 1: Four stat text boxes present (0.30 points)
    try:
        found_stats = []
        for stat in EXPECTED_STATS:
            matched = any(s['text'] == stat for s in shapes_info)
            found_stats.append((stat, matched))

        stats_found_count = sum(1 for _, m in found_stats if m)
        if stats_found_count == 4:
            print(f"PASS: Component 1 - All 4 stat values found on slide 3 (0.30 pts)")
            total_score += 0.30
        else:
            for stat, matched in found_stats:
                status = "found" if matched else "MISSING"
                print(f"  Stat '{stat}': {status}")
            print(f"FAIL: Component 1 - Only {stats_found_count}/4 stat values found")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Four subtitle text boxes present (0.20 points)
    try:
        found_subs = []
        for sub in EXPECTED_SUBTITLES:
            matched = any(s['text'] == sub for s in shapes_info)
            found_subs.append((sub, matched))

        subs_found_count = sum(1 for _, m in found_subs if m)
        if subs_found_count == 4:
            print(f"PASS: Component 2 - All 4 subtitle values found on slide 3 (0.20 pts)")
            total_score += 0.20
        else:
            for sub, matched in found_subs:
                status = "found" if matched else "MISSING"
                print(f"  Subtitle '{sub}': {status}")
            print(f"FAIL: Component 2 - Only {subs_found_count}/4 subtitles found")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Stat formatting - 48pt bold #2B6CB0 (0.20 points)
    try:
        stats_formatted = 0
        for stat in EXPECTED_STATS:
            for s in shapes_info:
                if s['text'] == stat:
                    size_ok = s['size'] is not None and abs(s['size'] - 609600) < 12700  # 48pt = 609600 EMU, allow ~1pt tolerance
                    bold_ok = s['bold'] is True
                    color_ok = s['color'] is not None and s['color'].upper() == '2B6CB0'
                    if size_ok and bold_ok and color_ok:
                        stats_formatted += 1
                    else:
                        print(f"  Stat '{stat}' formatting: size_ok={size_ok}(got {s['size']}), bold_ok={bold_ok}, color_ok={color_ok}(got {s['color']})")
                    break

        if stats_formatted == 4:
            print(f"PASS: Component 3 - All 4 stats correctly formatted: 48pt bold #2B6CB0 (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 3 - Only {stats_formatted}/4 stats correctly formatted")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Subtitle formatting - 16pt #666666 (0.15 points)
    try:
        subs_formatted = 0
        for sub in EXPECTED_SUBTITLES:
            for s in shapes_info:
                if s['text'] == sub:
                    size_ok = s['size'] is not None and abs(s['size'] - 203200) < 12700  # 16pt = 203200 EMU
                    color_ok = s['color'] is not None and s['color'].upper() == '666666'
                    if size_ok and color_ok:
                        subs_formatted += 1
                    else:
                        print(f"  Subtitle '{sub}' formatting: size_ok={size_ok}(got {s['size']}), color_ok={color_ok}(got {s['color']})")
                    break

        if subs_formatted == 4:
            print(f"PASS: Component 4 - All 4 subtitles correctly formatted: 16pt #666666 (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 - Only {subs_formatted}/4 subtitles correctly formatted")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Zoom entrance animations on slide 3 (0.15 points)
    # presetID="53" is Zoom, presetClass="entr" is entrance
    try:
        ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        ns = {'p': ns_p}

        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide3.xml') as f:
                root = ET.parse(f).getroot()

        # Find all animation nodes with presetID="53" and presetClass="entr"
        timing = root.find(f'.//{{{ns_p}}}timing')
        zoom_animations = []
        if timing is not None:
            for ctn in root.iter(f'{{{ns_p}}}cTn'):
                preset_id = ctn.get('presetID')
                preset_class = ctn.get('presetClass')
                if preset_id == '53' and preset_class == 'entr':
                    zoom_animations.append(ctn)

        # We expect at least 4 zoom entrance animations (one per stat, potentially paired with subtitle)
        # In the golden, there are 8 animations (4 stats + 4 subtitles)
        if len(zoom_animations) >= 4:
            print(f"PASS: Component 5 - Found {len(zoom_animations)} Zoom entrance animations on slide 3 (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 5 - Found only {len(zoom_animations)} Zoom entrance animations, need at least 4")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
