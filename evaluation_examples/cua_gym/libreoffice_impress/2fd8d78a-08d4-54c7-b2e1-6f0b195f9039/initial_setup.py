"""
Initial Setup: Design a team introduction section with 4 slides (slides 7-10)
Task ID: impress_sales_054
Domain: libreoffice_impress

Creates an 11-slide presentation 'Team_Deck.pptx' where slides 7-10 are blank.
Also creates 4 headshot placeholder images on the Desktop.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_054'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
DESKTOP = f'{WORKDIR}/Desktop'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_headshot_image(filepath, bg_color, initials):
    """Create a simple headshot placeholder image with colored background and initials."""
    size = 400
    img = Image.new('RGB', (size, size), bg_color)
    draw = ImageDraw.Draw(img)
    # Draw a circle for the face area
    margin = 20
    draw.ellipse([margin, margin, size - margin, size - margin],
                 fill=bg_color, outline='white', width=3)
    # Draw initials in center
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 80)
    except (IOError, OSError):
        font = ImageFont.load_default()
    bbox = draw.textbbox((0, 0), initials, font=font)
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]
    draw.text(((size - tw) / 2, (size - th) / 2 - 10), initials, fill='white', font=font)
    img.save(filepath)


def add_text_box(slide, left, top, width, height, text, font_size, bold=False, alignment=PP_ALIGN.LEFT, color=None):
    """Helper to add a styled text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    os.makedirs(DESKTOP, exist_ok=True)

    # Create headshot images on Desktop
    team_images = [
        ('team1.png', '#2E5090', 'SC'),  # Sarah Chen
        ('team2.png', '#1A7A4C', 'MJ'),  # Marcus Johnson
        ('team3.png', '#8B2252', 'LP'),  # Lisa Park
        ('team4.png', '#D4760A', 'DK'),  # David Kim
    ]
    for fname, color, initials in team_images:
        create_headshot_image(os.path.join(DESKTOP, fname), color, initials)
    print('Headshot images created on Desktop')

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0B, 0x2A, 0x52)
    add_text_box(slide1, Inches(1.5), Inches(1.5), Inches(10), Inches(1.5),
                 "NovaTech Solutions", 44, bold=True, alignment=PP_ALIGN.CENTER,
                 color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text_box(slide1, Inches(2), Inches(3.5), Inches(9), Inches(1),
                 "2025 Annual Sales Kickoff & Team Showcase", 24, alignment=PP_ALIGN.CENTER,
                 color=RGBColor(0xCC, 0xD5, 0xE0))
    add_text_box(slide1, Inches(3), Inches(5.0), Inches(7), Inches(0.6),
                 "March 15, 2025  |  San Francisco, CA", 16, alignment=PP_ALIGN.CENTER,
                 color=RGBColor(0x8F, 0xA8, 0xC8))

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide2, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                 "Agenda", 36, bold=True, color=RGBColor(0x0B, 0x2A, 0x52))
    agenda_items = [
        "1. Company Performance Overview",
        "2. Q4 2024 Revenue Results",
        "3. 2025 Growth Strategy",
        "4. Product Roadmap Updates",
        "5. Meet the Leadership Team",
        "6. Q&A and Networking",
    ]
    txBox = slide2.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(agenda_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.space_after = Pt(14)
        run = p.runs[0]
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 3: Company Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide3, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                 "Company Performance", 36, bold=True, color=RGBColor(0x0B, 0x2A, 0x52))
    metrics = [
        ("Revenue", "$127.4M", "+23% YoY"),
        ("Customers", "4,230", "+18% YoY"),
        ("Employee Count", "682", "+31% YoY"),
        ("NPS Score", "72", "+8 pts YoY"),
    ]
    for i, (label, value, change) in enumerate(metrics):
        x = Inches(0.8 + i * 3.0)
        add_text_box(slide3, x, Inches(2.0), Inches(2.5), Inches(0.6),
                     label, 16, color=RGBColor(0x66, 0x66, 0x66))
        add_text_box(slide3, x, Inches(2.6), Inches(2.5), Inches(0.8),
                     value, 32, bold=True, color=RGBColor(0x0B, 0x2A, 0x52))
        add_text_box(slide3, x, Inches(3.5), Inches(2.5), Inches(0.5),
                     change, 14, color=RGBColor(0x1A, 0x7A, 0x4C))

    # --- Slide 4: Revenue Breakdown ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide4, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                 "Revenue by Segment", 36, bold=True, color=RGBColor(0x0B, 0x2A, 0x52))
    # Table of revenue data
    table_shape = slide4.shapes.add_table(5, 4, Inches(1), Inches(2), Inches(10), Inches(3))
    table = table_shape.table
    headers = ["Segment", "Q4 2024", "Q4 2023", "Growth"]
    rows_data = [
        ["Enterprise", "$52.3M", "$41.8M", "+25.1%"],
        ["Mid-Market", "$38.7M", "$33.2M", "+16.6%"],
        ["SMB", "$24.1M", "$19.5M", "+23.6%"],
        ["Partner Channel", "$12.3M", "$8.9M", "+38.2%"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    for r, row in enumerate(rows_data, 1):
        for c, val in enumerate(row):
            table.cell(r, c).text = val

    # --- Slide 5: Growth Strategy ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide5, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                 "2025 Growth Strategy", 36, bold=True, color=RGBColor(0x0B, 0x2A, 0x52))
    strategies = [
        "Expand enterprise accounts in EMEA and APAC markets",
        "Launch AI-powered analytics module by Q2 2025",
        "Achieve SOC 2 Type II and ISO 27001 certification",
        "Grow partner ecosystem from 45 to 120 certified partners",
        "Increase net revenue retention to 135%",
    ]
    txBox5 = slide5.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10), Inches(4.5))
    tf5 = txBox5.text_frame
    tf5.word_wrap = True
    for i, item in enumerate(strategies):
        if i == 0:
            p = tf5.paragraphs[0]
        else:
            p = tf5.add_paragraph()
        p.text = item
        p.space_after = Pt(12)
        p.level = 0
        run = p.runs[0]
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 6: Product Roadmap ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide6, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                 "Product Roadmap Highlights", 36, bold=True, color=RGBColor(0x0B, 0x2A, 0x52))
    roadmap = [
        ("Q1 2025", "Enhanced reporting dashboard, API v3.0 release"),
        ("Q2 2025", "AI Analytics module, Mobile app refresh"),
        ("Q3 2025", "Enterprise SSO improvements, Workflow automation"),
        ("Q4 2025", "Platform v5.0, Advanced integrations suite"),
    ]
    for i, (quarter, desc) in enumerate(roadmap):
        y = Inches(2.0 + i * 1.2)
        add_text_box(slide6, Inches(1.2), y, Inches(2.5), Inches(0.6),
                     quarter, 20, bold=True, color=RGBColor(0x0B, 0x2A, 0x52))
        add_text_box(slide6, Inches(4.0), y, Inches(8), Inches(0.6),
                     desc, 18, color=RGBColor(0x33, 0x33, 0x33))

    # --- Slides 7-10: BLANK (Team Member Slides - to be completed by agent) ---
    for _ in range(4):
        prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # --- Slide 11: Closing / Contact ---
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    fill11 = slide11.background.fill
    fill11.solid()
    fill11.fore_color.rgb = RGBColor(0x0B, 0x2A, 0x52)
    add_text_box(slide11, Inches(2), Inches(2), Inches(9), Inches(1.5),
                 "Thank You!", 44, bold=True, alignment=PP_ALIGN.CENTER,
                 color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text_box(slide11, Inches(2), Inches(4.0), Inches(9), Inches(0.6),
                 "Questions? Reach us at team@novatech.io", 20, alignment=PP_ALIGN.CENTER,
                 color=RGBColor(0xCC, 0xD5, 0xE0))
    add_text_box(slide11, Inches(2), Inches(5.0), Inches(9), Inches(0.6),
                 "www.novatech-solutions.com  |  @NovatechHQ", 16, alignment=PP_ALIGN.CENTER,
                 color=RGBColor(0x8F, 0xA8, 0xC8))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
