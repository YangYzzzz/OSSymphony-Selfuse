"""
Reward Script: Team introduction slides (slides 7-10) with headshots, names, titles, bios
Task ID: impress_sales_054
Domain: libreoffice_impress
Scoring:
  Component 1: Slides 7-10 have shapes (not blank)           — 0.10
  Component 2: Oval shapes with image fill on each slide     — 0.25
  Component 3: Names correct + 24pt bold on each slide       — 0.25
  Component 4: Titles correct + 16pt on each slide           — 0.20
  Component 5: Bio text with 14pt + 2 lines on each slide    — 0.20
  Total: 1.0
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_054'

# Expected team members on slides 7-10 (0-indexed: 6-9)
TEAM_MEMBERS = [
    {"slide_idx": 6, "name": "Sarah Chen",      "title": "CEO",            "slide_num": 7},
    {"slide_idx": 7, "name": "Marcus Johnson",   "title": "CTO",            "slide_num": 8},
    {"slide_idx": 8, "name": "Lisa Park",        "title": "VP Sales",       "slide_num": 9},
    {"slide_idx": 9, "name": "David Kim",        "title": "VP Engineering", "slide_num": 10},
]


def persist_app_state():
    """Save any unsaved GUI state before verification."""
    try:
        os.environ["DISPLAY"] = ":0"
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        import time
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def get_all_text_from_slide(slide):
    """Extract all text shapes from a slide, including grouped shapes."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    results = []

    def extract(shape):
        if hasattr(shape, "text") and hasattr(shape, "text_frame"):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                extract(sub)

    for shape in slide.shapes:
        extract(shape)
    return results


def check_oval_with_image(slide):
    """Check if slide has an oval/ellipse shape with an image (blip) fill."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            xml_str = shape._element.xml
            # Check for ellipse geometry and blip image fill
            has_ellipse = 'prst="ellipse"' in xml_str or 'prst="oval"' in xml_str
            has_blip = 'blipFill' in xml_str or 'blip' in xml_str
            if has_ellipse and has_blip:
                return True
        # Also accept PICTURE shapes in oval-like scenarios
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            xml_str = shape._element.xml
            if 'ellipse' in xml_str or 'oval' in xml_str:
                return True
    return False


def find_text_matching(text_shapes, target_text, target_size_emu=None, target_bold=None):
    """Find a text shape containing the target text with optional font checks.
    Returns (found, size_match, bold_match)."""
    from pptx.util import Pt
    for shape in text_shapes:
        for para in shape.text_frame.paragraphs:
            full_text = para.text.strip()
            if target_text.lower() in full_text.lower():
                # Check font properties from first non-empty run
                runs = [r for r in para.runs if (r.text or "").strip()]
                if not runs:
                    return (True, target_size_emu is None, target_bold is None)
                run = runs[0]
                size_ok = (target_size_emu is None) or (
                    run.font.size is not None and run.font.size == target_size_emu
                )
                actual_bold = run.font.bold
                actual_bold = False if actual_bold is None else actual_bold
                bold_ok = (target_bold is None) or (actual_bold == target_bold)
                return (True, size_ok, bold_ok)
    return (False, False, False)


def count_bio_paragraphs(text_shapes, name, title):
    """Find the bio textbox (not name, not title) and count non-empty paragraphs.
    Also check if font size is 14pt (177800 EMU)."""
    from pptx.util import Pt
    target_bio_size = Pt(14)  # 177800 EMU
    for shape in text_shapes:
        all_text = shape.text_frame.text.strip()
        # Skip if it matches name or title
        if name.lower() in all_text.lower() and len(all_text) < len(name) + 10:
            continue
        if title.lower() in all_text.lower() and len(all_text) < len(title) + 10:
            continue
        # This could be the bio box: must have multiple paragraphs
        paras = [p for p in shape.text_frame.paragraphs if p.text.strip()]
        if len(paras) >= 2:
            # Check font size on first run of first para
            size_ok = any(
                r.font.size is not None and r.font.size == target_bio_size
                for p in paras
                for r in p.runs
                if (r.text or "").strip()
            )
            return (len(paras), size_ok)
    return (0, False)


def verify_task(file_path):
    """Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file must have at least 10 slides
    if len(prs.slides) < 10:
        print(f"FAIL: File has only {len(prs.slides)} slides, need at least 10")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slides 7-10 have shapes (not blank) (0.10 points)
    # Initial has 0 shapes on these slides; golden has 4 each
    try:
        slides_with_content = 0
        for member in TEAM_MEMBERS:
            slide = prs.slides[member["slide_idx"]]
            if len(slide.shapes) > 0:
                slides_with_content += 1
        if slides_with_content == 4:
            print(f"PASS: Component 1 — All 4 team slides have content (0.10 pts)")
            total_score += 0.10
        elif slides_with_content > 0:
            partial = round(0.10 * slides_with_content / 4, 4)
            print(f"PARTIAL: Component 1 — {slides_with_content}/4 slides have content ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — All team slides are blank")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Oval shapes with image fill on each slide (0.25 points)
    try:
        ovals_found = 0
        for member in TEAM_MEMBERS:
            slide = prs.slides[member["slide_idx"]]
            if check_oval_with_image(slide):
                ovals_found += 1
                print(f"  Slide {member['slide_num']}: oval with image found")
            else:
                print(f"  Slide {member['slide_num']}: no oval with image")
        if ovals_found == 4:
            print(f"PASS: Component 2 — All 4 slides have oval image shapes (0.25 pts)")
            total_score += 0.25
        elif ovals_found > 0:
            partial = round(0.25 * ovals_found / 4, 4)
            print(f"PARTIAL: Component 2 — {ovals_found}/4 oval images ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No oval image shapes found on team slides")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Names correct + 24pt bold on each slide (0.25 points)
    try:
        target_name_size = Pt(24)  # 304800 EMU
        names_correct = 0
        for member in TEAM_MEMBERS:
            slide = prs.slides[member["slide_idx"]]
            text_shapes = get_all_text_from_slide(slide)
            found, size_ok, bold_ok = find_text_matching(
                text_shapes, member["name"], target_name_size, True
            )
            if found and size_ok and bold_ok:
                names_correct += 1
                print(f"  Slide {member['slide_num']}: name '{member['name']}' correct (24pt bold)")
            elif found:
                # Give half credit if name present but wrong formatting
                names_correct += 0.5
                print(f"  Slide {member['slide_num']}: name '{member['name']}' found but formatting off (size_ok={size_ok}, bold_ok={bold_ok})")
            else:
                print(f"  Slide {member['slide_num']}: name '{member['name']}' NOT found")
        if names_correct >= 4:
            print(f"PASS: Component 3 — All 4 names correct with 24pt bold (0.25 pts)")
            total_score += 0.25
        elif names_correct > 0:
            partial = round(0.25 * names_correct / 4, 4)
            print(f"PARTIAL: Component 3 — {names_correct}/4 names correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No correct names found on team slides")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Titles correct + 16pt on each slide (0.20 points)
    try:
        target_title_size = Pt(16)  # 203200 EMU
        titles_correct = 0
        for member in TEAM_MEMBERS:
            slide = prs.slides[member["slide_idx"]]
            text_shapes = get_all_text_from_slide(slide)
            found, size_ok, _ = find_text_matching(
                text_shapes, member["title"], target_title_size, None
            )
            if found and size_ok:
                titles_correct += 1
                print(f"  Slide {member['slide_num']}: title '{member['title']}' correct (16pt)")
            elif found:
                titles_correct += 0.5
                print(f"  Slide {member['slide_num']}: title '{member['title']}' found but size wrong")
            else:
                print(f"  Slide {member['slide_num']}: title '{member['title']}' NOT found")
        if titles_correct >= 4:
            print(f"PASS: Component 4 — All 4 titles correct with 16pt (0.20 pts)")
            total_score += 0.20
        elif titles_correct > 0:
            partial = round(0.20 * titles_correct / 4, 4)
            print(f"PARTIAL: Component 4 — {titles_correct}/4 titles correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — No correct titles found on team slides")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Bio text with 14pt + 2 lines on each slide (0.20 points)
    try:
        bios_correct = 0
        for member in TEAM_MEMBERS:
            slide = prs.slides[member["slide_idx"]]
            text_shapes = get_all_text_from_slide(slide)
            para_count, size_ok = count_bio_paragraphs(
                text_shapes, member["name"], member["title"]
            )
            if para_count >= 2 and size_ok:
                bios_correct += 1
                print(f"  Slide {member['slide_num']}: bio has {para_count} lines at 14pt")
            elif para_count >= 2:
                bios_correct += 0.5
                print(f"  Slide {member['slide_num']}: bio has {para_count} lines but size wrong")
            else:
                print(f"  Slide {member['slide_num']}: bio insufficient (paras={para_count})")
        if bios_correct >= 4:
            print(f"PASS: Component 5 — All 4 bios correct with 14pt + 2 lines (0.20 pts)")
            total_score += 0.20
        elif bios_correct > 0:
            partial = round(0.20 * bios_correct / 4, 4)
            print(f"PARTIAL: Component 5 — {bios_correct}/4 bios correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 — No correct bios found on team slides")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
