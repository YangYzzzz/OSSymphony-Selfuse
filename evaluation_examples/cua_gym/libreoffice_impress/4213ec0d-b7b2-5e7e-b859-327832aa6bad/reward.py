"""
Reward Script: Interactive menu-driven Company Handbook presentation
Task ID: impress_wf_015
Domain: libreoffice_impress
Scoring:
  C1 (0.15): 5 slides total
  C2 (0.25): Slide 1 title 'Company Handbook' + 4 labelled menu shapes
  C3 (0.20): Menu shapes have #2196F3 fill, rounded corners, drop shadow
  C4 (0.20): Menu shapes hyperlink to slides 2-5 respectively
  C5 (0.20): Slides 2-5 have 'Back to Menu' text with hyperlink to slide 1
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_015'

# XML namespaces used for low-level checks
NS = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}


def get_slide_rels(pptx_path, slide_num):
    """Return dict of rId -> target filename for a slide's relationships."""
    rels = {}
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            rels_path = f'ppt/slides/_rels/slide{slide_num}.xml.rels'
            with zf.open(rels_path) as f:
                root = ET.parse(f).getroot()
                for rel in root:
                    rels[rel.get('Id')] = rel.get('Target', '')
    except Exception:
        pass
    return rels


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.oxml.ns import qn
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # ----------------------------------------------------------------
    # Component 1: Presentation has exactly 5 slides (0.15 points)
    # ----------------------------------------------------------------
    try:
        slide_count = len(prs.slides)
        if slide_count == 5:
            print(f"PASS: Component 1 — slide count is 5 (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — expected 5 slides, found {slide_count}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Early exit if not enough slides for remaining checks
    if len(prs.slides) < 5:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # ----------------------------------------------------------------
    # Component 2: Slide 1 has title 'Company Handbook' + 4 menu
    #              rectangles labelled Mission, Policies, Benefits, FAQ
    #              (0.25 points)
    # ----------------------------------------------------------------
    try:
        slide1 = prs.slides[0]
        expected_labels = {'mission', 'policies', 'benefits', 'faq'}

        # Check title
        title_found = False
        for shape in slide1.shapes:
            if hasattr(shape, 'text') and 'company handbook' in (shape.text or '').lower():
                title_found = True
                break

        # Collect menu rectangle labels (AUTO_SHAPE type == 1)
        menu_labels = set()
        for shape in slide1.shapes:
            if shape.shape_type == 1:  # AUTO_SHAPE
                txt = (shape.text or '').strip().lower()
                if txt in expected_labels:
                    menu_labels.add(txt)

        if title_found and menu_labels == expected_labels:
            print(f"PASS: Component 2 — title 'Company Handbook' found + 4 menu labels {menu_labels} (0.25 pts)")
            total_score += 0.25
        else:
            missing = expected_labels - menu_labels
            print(f"FAIL: Component 2 — title_found={title_found}, menu_labels={menu_labels}, missing={missing}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # ----------------------------------------------------------------
    # Component 3: Menu shapes have blue fill (#2196F3), rounded
    #              corners (roundRect), and drop shadow (0.20 points)
    # ----------------------------------------------------------------
    try:
        from pptx.oxml.ns import qn
        slide1 = prs.slides[0]
        styled_count = 0

        for shape in slide1.shapes:
            if shape.shape_type != 1:  # only AUTO_SHAPE
                continue
            txt = (shape.text or '').strip().lower()
            if txt not in expected_labels:
                continue

            sp = shape._element

            # Check fill color
            solid_fill = sp.find('.//' + qn('a:solidFill'))
            fill_ok = False
            if solid_fill is not None:
                srgb = solid_fill.find(qn('a:srgbClr'))
                if srgb is not None and srgb.get('val', '').upper() == '2196F3':
                    fill_ok = True

            # Check rounded corners (prstGeom == roundRect)
            prst_geom = sp.find('.//' + qn('a:prstGeom'))
            round_ok = prst_geom is not None and prst_geom.get('prst') == 'roundRect'

            # Check drop shadow
            effect_lst = sp.find('.//' + qn('a:effectLst'))
            shadow_ok = False
            if effect_lst is not None:
                shadow_ok = effect_lst.find(qn('a:outerShdw')) is not None

            if fill_ok and round_ok and shadow_ok:
                styled_count += 1
            else:
                print(f"  Shape '{txt}': fill_ok={fill_ok}, round_ok={round_ok}, shadow_ok={shadow_ok}")

        if styled_count == 4:
            print(f"PASS: Component 3 — all 4 menu shapes have #2196F3, roundRect, drop shadow (0.20 pts)")
            total_score += 0.20
        elif styled_count > 0:
            partial = round(0.20 * styled_count / 4, 2)
            print(f"PARTIAL: Component 3 — {styled_count}/4 shapes styled correctly ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — no menu shapes have correct styling")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # ----------------------------------------------------------------
    # Component 4: Menu shapes hyperlink to slides 2-5 respectively
    #              (0.20 points)
    # ----------------------------------------------------------------
    try:
        from pptx.oxml.ns import qn
        slide1 = prs.slides[0]
        rels = get_slide_rels(file_path, 1)

        # Map expected label -> target slide file
        label_to_target = {
            'mission': 'slide2.xml',
            'policies': 'slide3.xml',
            'benefits': 'slide4.xml',
            'faq': 'slide5.xml',
        }

        correct_links = 0
        for shape in slide1.shapes:
            if shape.shape_type != 1:
                continue
            txt = (shape.text or '').strip().lower()
            if txt not in label_to_target:
                continue

            expected_target = label_to_target[txt]
            sp = shape._element

            # Look for shape-level hyperlink (hlinkClick on cNvPr)
            cNvPr = sp.find('.//' + qn('p:cNvPr'))
            if cNvPr is not None:
                hlink = cNvPr.find(qn('a:hlinkClick'))
                if hlink is not None:
                    action = hlink.get('action', '')
                    rId = hlink.get(qn('r:id'), '')
                    target = rels.get(rId, '')
                    if 'hlinksldjump' in action and target == expected_target:
                        correct_links += 1
                        continue

            # Also check run-level hyperlinks as fallback
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        rPr = run._r.find(qn('a:rPr'))
                        if rPr is not None:
                            hlink = rPr.find(qn('a:hlinkClick'))
                            if hlink is not None:
                                action = hlink.get('action', '')
                                rId = hlink.get(qn('r:id'), '')
                                target = rels.get(rId, '')
                                if 'hlinksldjump' in action and target == expected_target:
                                    correct_links += 1
                                    break

        if correct_links == 4:
            print(f"PASS: Component 4 — all 4 menu shapes link to correct slides (0.20 pts)")
            total_score += 0.20
        elif correct_links > 0:
            partial = round(0.20 * correct_links / 4, 2)
            print(f"PARTIAL: Component 4 — {correct_links}/4 correct hyperlinks ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — no menu shapes have correct hyperlinks")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # ----------------------------------------------------------------
    # Component 5: Slides 2-5 each have 'Back to Menu' text with
    #              hyperlink to slide 1 (0.20 points)
    # ----------------------------------------------------------------
    try:
        from pptx.oxml.ns import qn
        back_links_ok = 0

        for slide_idx in range(1, 5):  # slides 2-5 (0-indexed: 1-4)
            slide = prs.slides[slide_idx]
            slide_num = slide_idx + 1
            rels = get_slide_rels(file_path, slide_num)

            found_back = False
            for shape in slide.shapes:
                if not hasattr(shape, 'text'):
                    continue
                if 'back to menu' not in (shape.text or '').lower():
                    continue

                # Check for hyperlink in runs
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            rPr = run._r.find(qn('a:rPr'))
                            if rPr is not None:
                                hlink = rPr.find(qn('a:hlinkClick'))
                                if hlink is not None:
                                    action = hlink.get('action', '')
                                    rId = hlink.get(qn('r:id'), '')
                                    target = rels.get(rId, '')
                                    if 'hlinksldjump' in action and target == 'slide1.xml':
                                        found_back = True
                                        break
                        if found_back:
                            break

                # Also check shape-level hyperlink
                if not found_back:
                    sp = shape._element
                    cNvPr = sp.find('.//' + qn('p:cNvPr'))
                    if cNvPr is not None:
                        hlink = cNvPr.find(qn('a:hlinkClick'))
                        if hlink is not None:
                            action = hlink.get('action', '')
                            rId = hlink.get(qn('r:id'), '')
                            target = rels.get(rId, '')
                            if 'hlinksldjump' in action and target == 'slide1.xml':
                                found_back = True

                if found_back:
                    break

            if found_back:
                back_links_ok += 1
            else:
                print(f"  Slide {slide_num}: 'Back to Menu' link to slide 1 NOT found")

        if back_links_ok == 4:
            print(f"PASS: Component 5 — all 4 content slides have 'Back to Menu' -> slide 1 (0.20 pts)")
            total_score += 0.20
        elif back_links_ok > 0:
            partial = round(0.20 * back_links_ok / 4, 2)
            print(f"PARTIAL: Component 5 — {back_links_ok}/4 content slides have correct back link ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 — no content slides have 'Back to Menu' linking to slide 1")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
