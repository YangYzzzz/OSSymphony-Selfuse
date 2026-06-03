"""
Reward Script: Conference Poster Design in LibreOffice Impress
Task ID: impress_wf_014
Domain: libreoffice_impress
Scoring:
  - Slide dimensions 48x36 inches (0.15)
  - Single slide (0.05)
  - Title with correct text, color #003366, bold (0.20)
  - Three-column sections: Introduction, Methods, Results+Conclusions (0.20)
  - Chart placeholder rectangle (0.10)
  - References section at bottom (0.10)
  - Horizontal decorative lines (0.10)
  - Section headers in #003366 (0.10)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_014'
FILE_PATH = os.path.join(WORKDIR, 'Desktop', 'Conference_Poster.pptx')


def get_all_text_shapes(slide):
    """Recursively get all shapes with text frames, including grouped shapes."""
    def extract(shape):
        results = []
        if hasattr(shape, "text") and hasattr(shape, "text_frame"):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.extend(extract(sub))
        return results
    out = []
    for shape in slide.shapes:
        out.extend(extract(shape))
    return out


def get_font_color_hex(run):
    """Safely extract font color as hex string."""
    try:
        if run.font.color and run.font.color.type is not None:
            return str(run.font.color.rgb).upper()
    except Exception:
        pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Single slide (0.05 points)
    try:
        num_slides = len(prs.slides)
        if num_slides == 1:
            print(f"PASS: Component 1 -- Single slide confirmed (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 1 -- Expected 1 slide, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Slide dimensions 48x36 inches (0.15 points)
    try:
        width_inches = prs.slide_width / 914400
        height_inches = prs.slide_height / 914400
        width_ok = abs(width_inches - 48.0) < 0.5
        height_ok = abs(height_inches - 36.0) < 0.5
        if width_ok and height_ok:
            print(f"PASS: Component 2 -- Slide dimensions {width_inches:.1f}x{height_inches:.1f} inches (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 -- Expected 48x36 inches, found {width_inches:.1f}x{height_inches:.1f}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Get slide for remaining checks
    if len(prs.slides) < 1:
        print("CRITICAL: No slides found, cannot proceed")
        print(f"REWARD: {total_score}")
        return total_score

    slide = prs.slides[0]
    text_shapes = get_all_text_shapes(slide)

    # Component 3: Title text "Machine Learning for Climate Prediction" with #003366 and bold (0.20 points)
    try:
        title_found = False
        title_color_ok = False
        title_bold_ok = False
        for shape in text_shapes:
            for para in shape.text_frame.paragraphs:
                if 'machine learning for climate prediction' in para.text.lower():
                    title_found = True
                    for run in para.runs:
                        if 'machine learning' in run.text.lower():
                            color = get_font_color_hex(run)
                            if color == '003366':
                                title_color_ok = True
                            bold = run.font.bold
                            if bold is True:
                                title_bold_ok = True
                    break
            if title_found:
                break

        sub_score = 0.0
        if title_found:
            sub_score += 0.10
        if title_color_ok:
            sub_score += 0.05
        if title_bold_ok:
            sub_score += 0.05

        if sub_score > 0:
            print(f"PASS: Component 3 -- Title: found={title_found}, color=#003366={title_color_ok}, bold={title_bold_ok} ({sub_score} pts)")
            total_score += sub_score
        else:
            print(f"FAIL: Component 3 -- Title 'Machine Learning for Climate Prediction' not found")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Three-column content sections (Introduction, Methods, Results+Conclusions) (0.20 points)
    try:
        sections_found = {'introduction': False, 'methods': False, 'results': False, 'conclusions': False}
        for shape in text_shapes:
            full_text = shape.text_frame.text.lower()
            for section in sections_found:
                if section in full_text:
                    # Check it appears as a header (first paragraph or prominent)
                    for para in shape.text_frame.paragraphs:
                        if section in para.text.lower().strip().lower():
                            sections_found[section] = True
                            break

        found_count = sum(1 for v in sections_found.values() if v)
        sub_score = found_count * 0.05  # 0.05 per section, 4 sections = 0.20
        if sub_score > 0:
            details = ", ".join(f"{k}={'YES' if v else 'NO'}" for k, v in sections_found.items())
            print(f"PASS: Component 4 -- Sections: {details} ({sub_score} pts)")
            total_score += sub_score
        else:
            print(f"FAIL: Component 4 -- No content sections found")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Chart placeholder rectangle in Results area (0.10 points)
    try:
        chart_placeholder_found = False
        for shape in slide.shapes:
            # Look for an AutoShape (rectangle) with chart-related text
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if hasattr(shape, 'text_frame'):
                    shape_text = shape.text_frame.text.lower()
                    if 'chart' in shape_text or 'placeholder' in shape_text:
                        chart_placeholder_found = True
                        break
                # Also check for rectangle shapes in the Results column area (right third)
                # that have significant height (not a line)
                if shape.height > 914400:  # taller than 1 inch (not a decorative line)
                    # Check if it's in the right portion of the slide (Results column area)
                    if shape.left > prs.slide_width * 0.5:
                        chart_placeholder_found = True
                        break

        if chart_placeholder_found:
            print(f"PASS: Component 5 -- Chart placeholder rectangle found (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 5 -- No chart placeholder rectangle found in Results area")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # Component 6: References section at bottom (0.10 points)
    try:
        references_found = False
        for shape in text_shapes:
            full_text = shape.text_frame.text.lower()
            if 'references' in full_text:
                # Check it's in the bottom portion of the slide (last 30%)
                shape_top_ratio = shape.top / prs.slide_height if prs.slide_height > 0 else 0
                if shape_top_ratio > 0.6:
                    references_found = True
                    break

        if references_found:
            print(f"PASS: Component 6 -- References section at bottom (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 6 -- References section not found at bottom of slide")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    # Component 7: Horizontal decorative lines (thin rectangles) (0.10 points)
    try:
        line_count = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # Decorative lines are very thin rectangles spanning significant width
                if shape.height < 914400 * 0.5 and shape.width > prs.slide_width * 0.3:
                    line_count += 1

        if line_count >= 2:
            print(f"PASS: Component 7 -- Found {line_count} horizontal decorative lines (0.10 pts)")
            total_score += 0.10
        elif line_count == 1:
            print(f"PARTIAL: Component 7 -- Found {line_count} horizontal line (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 7 -- No horizontal decorative lines found")
    except Exception as e:
        print(f"ERROR: Component 7 -- {e}")

    # Component 8: Section headers in #003366 color (0.10 points)
    try:
        header_keywords = ['introduction', 'methods', 'results', 'conclusions', 'references']
        headers_with_correct_color = 0
        total_headers_found = 0

        for shape in text_shapes:
            for para in shape.text_frame.paragraphs:
                para_text = para.text.strip().lower()
                if para_text in header_keywords:
                    total_headers_found += 1
                    for run in para.runs:
                        color = get_font_color_hex(run)
                        if color == '003366':
                            headers_with_correct_color += 1
                            break

        if total_headers_found >= 3 and headers_with_correct_color >= 3:
            print(f"PASS: Component 8 -- {headers_with_correct_color}/{total_headers_found} section headers in #003366 (0.10 pts)")
            total_score += 0.10
        elif headers_with_correct_color >= 1:
            sub = round(0.10 * headers_with_correct_color / max(total_headers_found, 3), 2)
            print(f"PARTIAL: Component 8 -- {headers_with_correct_color}/{total_headers_found} headers in #003366 ({sub} pts)")
            total_score += sub
        else:
            print(f"FAIL: Component 8 -- No section headers found with #003366 color (found {total_headers_found} headers)")
    except Exception as e:
        print(f"ERROR: Component 8 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(FILE_PATH)
