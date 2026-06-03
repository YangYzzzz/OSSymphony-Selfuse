"""
initial_setup.py - Creates Anthropology_Talk.pptx with 10 slides, no transitions.
"""
import os
import subprocess
import time

subprocess.run(['pip3', 'install', 'python-pptx'], check=True,
               stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
OUTPUT_PATH = os.path.join(WORKDIR, 'Anthropology_Talk.pptx')

# Slide content: (title, bullet_points)
SLIDES = [
    ("Anthropology Talk", [
        "An Introduction to Cultural and Biological Anthropology",
        "Presented by Dr. Sarah Chen",
        "Department of Anthropology, Fall 2025"
    ]),
    ("What is Anthropology?", [
        "The study of humans, past and present",
        "Combines biological, social, and cultural perspectives",
        "Four-field approach: cultural, biological, linguistic, archaeological",
        "Seeks to understand human diversity and universals"
    ]),
    ("Cultural Anthropology", [
        "Study of cultural variation among humans",
        "Ethnography as primary research method",
        "Key concepts: ethnocentrism, cultural relativism",
        "Examines kinship, religion, economics, and politics",
        "Fieldwork traditions from Malinowski to modern participatory methods"
    ]),
    ("Biological Anthropology", [
        "Study of human biological evolution and variation",
        "Paleoanthropology: fossil record of human ancestors",
        "Primatology: behavior of non-human primates",
        "Forensic anthropology: skeletal identification",
        "Human genetics and population biology"
    ]),
    ("Linguistic Anthropology", [
        "Relationship between language and culture",
        "Sapir-Whorf hypothesis and linguistic relativity",
        "Language documentation and endangered languages",
        "Sociolinguistics and code-switching",
        "Historical linguistics and language families"
    ]),
    ("Archaeological Anthropology", [
        "Recovering and interpreting material remains",
        "Stratigraphy and dating methods",
        "Excavation techniques and site preservation",
        "Artifacts, ecofacts, and features",
        "Ethnoarchaeology and experimental archaeology"
    ]),
    ("Human Evolution", [
        "From Australopithecus to Homo sapiens",
        "Bipedalism as a defining adaptation",
        "Tool use and the development of culture",
        "Out of Africa theory vs. multiregional hypothesis",
        "Recent discoveries: Homo naledi, Denisovans"
    ]),
    ("Ethnographic Methods", [
        "Participant observation in field settings",
        "Semi-structured and open-ended interviews",
        "Life histories and narrative analysis",
        "Visual anthropology and digital ethnography",
        "Ethical considerations in fieldwork"
    ]),
    ("Applied Anthropology", [
        "Using anthropological knowledge to solve practical problems",
        "Medical anthropology and global health",
        "Development anthropology and policy making",
        "Corporate ethnography and UX research",
        "Cultural resource management"
    ]),
    ("Conclusion and Future Directions", [
        "Anthropology in the age of globalization",
        "Digital anthropology and virtual communities",
        "Climate change and environmental anthropology",
        "Decolonizing anthropological practice",
        "Questions and discussion"
    ]),
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

for i, (title, bullets) in enumerate(SLIDES):
    if i == 0:
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
        slide.shapes.title.text = title
        if slide.placeholders[1]:
            slide.placeholders[1].text = bullets[0]
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for j, bullet in enumerate(bullets):
            if j == 0:
                tf.paragraphs[0].text = bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

prs.save(OUTPUT_PATH)
print(f"Created {OUTPUT_PATH} with {len(prs.slides)} slides, no transitions.")

# Open in LibreOffice Impress
env = os.environ.copy()
env['DISPLAY'] = ':0'
subprocess.Popen(
    ['libreoffice', '--impress', OUTPUT_PATH],
    stdout=subprocess.DEVNULL,
    stderr=subprocess.DEVNULL,
    env=env,
)
time.sleep(2)
print("LibreOffice Impress launched.")
