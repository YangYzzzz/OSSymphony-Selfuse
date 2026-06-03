"""
Reward Script: Create a 5x5 multiplication table on slide 2
Task ID: impress_tct_027
Domain: libreoffice_impress
Scoring:
  Component 1 (0.2): Table exists on slide 2 with 6x6 dimensions
  Component 2 (0.2): Header row contains '1'-'5', cell(0,0) is empty
  Component 3 (0.2): Header column contains '1'-'5'
  Component 4 (0.4): Body cells contain correct products (partial credit)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_027'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 2")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]  # 0-indexed, slide 2

    # Find the table on slide 2
    table = None
    for shape in slide2.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            break

    # Component 1: Table exists on slide 2 with 6x6 dimensions (0.2 points)
    try:
        if table is None:
            print("FAIL: Component 1 -- No table found on slide 2")
        elif len(table.rows) == 6 and len(table.columns) == 6:
            print(f"PASS: Component 1 -- 6x6 table found on slide 2 (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 1 -- Table dimensions are {len(table.rows)}x{len(table.columns)}, expected 6x6")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # If no table found, cannot check further components
    if table is None:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: Header row (row 0) -- cell(0,0) empty, cols 1-5 contain '1'-'5' (0.2 points)
    try:
        cell_00 = table.cell(0, 0).text.strip()
        header_row_ok = (cell_00 == '')
        if not header_row_ok:
            print(f"FAIL: Component 2 -- cell(0,0) should be empty, found '{cell_00}'")

        header_row_values = []
        header_row_correct = 0
        for c in range(1, min(6, len(table.columns))):
            val = table.cell(0, c).text.strip()
            header_row_values.append(val)
            if val == str(c):
                header_row_correct += 1

        if header_row_ok and header_row_correct == 5:
            print(f"PASS: Component 2 -- Header row correct: cell(0,0)='' and row 0 cols 1-5 = {header_row_values} (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 2 -- Header row: cell(0,0)='{cell_00}', values={header_row_values}, {header_row_correct}/5 correct")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Header column (col 0, rows 1-5) contain '1'-'5' (0.2 points)
    try:
        header_col_values = []
        header_col_correct = 0
        for r in range(1, min(6, len(table.rows))):
            val = table.cell(r, 0).text.strip()
            header_col_values.append(val)
            if val == str(r):
                header_col_correct += 1

        if header_col_correct == 5:
            print(f"PASS: Component 3 -- Header column correct: {header_col_values} (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 3 -- Header column: {header_col_values}, {header_col_correct}/5 correct")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Body cells contain correct products (0.4 points, partial credit)
    # 25 body cells (5x5), each worth 0.4/25 = 0.016 points
    try:
        correct_cells = 0
        total_body_cells = 25
        wrong_cells = []
        for r in range(1, min(6, len(table.rows))):
            for c in range(1, min(6, len(table.columns))):
                expected = str(r * c)
                actual = table.cell(r, c).text.strip()
                if actual == expected:
                    correct_cells += 1
                else:
                    wrong_cells.append(f"cell({r},{c}): expected '{expected}', got '{actual}'")

        points = 0.4 * (correct_cells / total_body_cells)
        if correct_cells == total_body_cells:
            print(f"PASS: Component 4 -- All {total_body_cells} body cells correct (0.4 pts)")
        else:
            print(f"PARTIAL: Component 4 -- {correct_cells}/{total_body_cells} body cells correct ({points:.3f} pts)")
            for w in wrong_cells[:5]:  # Show up to 5 wrong cells
                print(f"  {w}")

        if correct_cells > 0:
            total_score += round(points, 4)
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
