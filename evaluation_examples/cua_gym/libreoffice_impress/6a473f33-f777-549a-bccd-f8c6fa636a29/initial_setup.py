"""
Initial Setup: Create Math_Lesson presentation with 4 slides; slide 2 has title but no table.
Task ID: impress_tct_027
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_027'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Math Lesson"
    slide1.placeholders[1].text = "An Interactive Guide to Basic Arithmetic"

    # --- Slide 2: Multiplication Table (title only, NO table yet) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Multiplication Table"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # --- Slide 3: Addition Practice ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Addition Practice"
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.runs[0]
    r3.font.name = "Arial"
    r3.font.size = Pt(32)
    r3.font.bold = True
    r3.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # Add some addition problems
    content3 = slide3.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    ctf3 = content3.text_frame
    ctf3.word_wrap = True
    problems = [
        "1)  14 + 27 = ___",
        "2)  38 + 56 = ___",
        "3)  123 + 489 = ___",
        "4)  67 + 245 = ___",
        "5)  1,024 + 3,976 = ___",
    ]
    for i, prob in enumerate(problems):
        if i == 0:
            cp = ctf3.paragraphs[0]
        else:
            cp = ctf3.add_paragraph()
        cp.text = prob
        cp.space_after = Pt(14)
        cr = cp.runs[0]
        cr.font.name = "Courier New"
        cr.font.size = Pt(20)

    # --- Slide 4: Division Basics ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Division Basics"
    p4.alignment = PP_ALIGN.CENTER
    r4 = p4.runs[0]
    r4.font.name = "Arial"
    r4.font.size = Pt(32)
    r4.font.bold = True
    r4.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    content4 = slide4.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    ctf4 = content4.text_frame
    ctf4.word_wrap = True
    div_content = [
        "Division is the inverse of multiplication.",
        "",
        "Key concepts:",
        "  - Dividend / Divisor = Quotient",
        "  - 20 / 4 = 5",
        "  - 36 / 6 = 6",
        "  - 100 / 25 = 4",
    ]
    for i, line in enumerate(div_content):
        if i == 0:
            dp = ctf4.paragraphs[0]
        else:
            dp = ctf4.add_paragraph()
        dp.text = line
        if line and dp.runs:
            dr = dp.runs[0]
            dr.font.name = "Arial"
            dr.font.size = Pt(18)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
