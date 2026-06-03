"""
Initial Setup: Marketing Lecture presentation with 7 slides. Slide 3 has 5 bullet points with no animations.
Task ID: impress_teach_023
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_023'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullets, layout_idx=1):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
    return slide


def add_blank_title_only(prs, title_text, body_text):
    """Add a slide with title and a text box for body content."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title text box
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    # Body text box
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = body_text
    p2.font.size = Pt(18)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    add_title_slide(prs, "Digital Marketing Strategy 2025",
                    "Quarterly Review & Forward Planning\nMarketing Department | March 2025")

    # Slide 2: Agenda
    add_content_slide(prs, "Agenda", [
        "Q4 2024 Performance Review",
        "Market Landscape & Competitor Analysis",
        "Customer Segmentation Insights",
        "Channel Strategy for Q1-Q2 2025",
        "Budget Allocation & KPIs",
        "Team Assignments & Next Steps",
    ])

    # Slide 3: Customer Segmentation Insights (THE TARGET SLIDE - 5 bullets, no animations)
    add_content_slide(prs, "Customer Segmentation Insights", [
        "Enterprise accounts (500+ employees) represent 42% of revenue with highest retention at 94%",
        "Mid-market segment grew 23% YoY, driven by product-led growth initiatives",
        "SMB cohort shows strongest response to email nurture campaigns with 18% conversion rate",
        "Geographic expansion into APAC contributed 15% of new logo acquisition in Q4",
        "Customer lifetime value increased to $34,200 across all segments, up 12% from prior year",
    ])

    # Slide 4: Channel Strategy
    add_content_slide(prs, "Channel Strategy for Q1-Q2 2025", [
        "Paid Search: Increase budget by 15% focused on high-intent keywords",
        "Content Marketing: Launch weekly thought leadership series on LinkedIn",
        "Email: Implement AI-driven personalization for nurture sequences",
        "Events: Sponsor 3 industry conferences (SaaStr, Dreamforce, Web Summit)",
        "Partnerships: Develop co-marketing programs with top 5 integration partners",
    ])

    # Slide 5: Budget Allocation
    add_content_slide(prs, "Budget Allocation", [
        "Total Q1-Q2 Budget: $2.4M (8% increase from prior period)",
        "Digital Advertising: $720K (30%)",
        "Content & Creative: $480K (20%)",
        "Events & Sponsorships: $360K (15%)",
        "Technology & Tools: $840K (35%)",
    ])

    # Slide 6: KPIs & Targets
    add_content_slide(prs, "Key Performance Indicators", [
        "Website Traffic: 150K monthly unique visitors (currently 112K)",
        "Marketing Qualified Leads: 800/month target (currently 620)",
        "Cost Per Acquisition: Reduce to $185 from current $224",
        "Pipeline Contribution: $8.5M in marketing-sourced pipeline",
    ])

    # Slide 7: Next Steps
    add_content_slide(prs, "Next Steps & Action Items", [
        "Finalize creative briefs for Q1 campaigns by March 28",
        "Complete vendor selection for marketing automation upgrade",
        "Schedule cross-functional alignment with Sales and Product teams",
        "Launch customer advisory board recruitment by April 15",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
