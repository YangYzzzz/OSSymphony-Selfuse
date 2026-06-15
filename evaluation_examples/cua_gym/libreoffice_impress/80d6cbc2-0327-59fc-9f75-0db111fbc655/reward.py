"""
Reward Script: Set text box internal margins on slide 6
Task ID: impress_tct_085
Domain: libreoffice_impress
Scoring: 4 components (top/bottom/left/right margins) x 0.25 each = 1.0
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_085'

# Expected margin values in EMU (914400 EMU = 1 inch)
EXPECTED_TOP = Inches(0.2)      # 182880 EMU
EXPECTED_BOTTOM = Inches(0.2)   # 182880 EMU
EXPECTED_LEFT = Inches(0.3)     # 274320 EMU
EXPECTED_RIGHT = Inches(0.3)    # 274320 EMU

# Tolerance: 2% relative or 5000 EMU absolute (whichever is larger)
def is_close(actual, expected, rel_tol=0.02, abs_tol=5000):
    """Check if two EMU values are approximately equal."""
    if actual is None:
        return False
    diff = abs(actual - expected)
    threshold = max(abs(expected * rel_tol), abs_tol)
    return diff <= threshold


def persist_app_state(domain):
    """Save any unsaved GUI state before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify that the text box on slide 6 has the correct internal margins.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[5]  # Slide 6 (0-indexed)

    # Find the text box by name or as the non-placeholder shape with text
    # Note: shape_type may vary (TEXT_BOX=17 or AUTO_SHAPE=1) depending on how it was created
    textbox = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "TextBox 2":
            textbox = shape
            break
    # Fallback: find any non-placeholder shape with text
    if textbox is None:
        for shape in slide.shapes:
            if shape.has_text_frame and int(shape.shape_type) != 14:  # not PLACEHOLDER
                textbox = shape
                break

    if textbox is None:
        print("FAIL: No text box found on slide 6")
        print("REWARD: 0.0")
        return 0.0

    tf = textbox.text_frame
    print(f"INFO: Found text box '{textbox.name}' on slide 6")

    # Component 1: Top margin = 0.2 inches (0.25 points)
    try:
        actual_top = tf.margin_top
        if is_close(actual_top, EXPECTED_TOP):
            print(f"PASS: Component 1 -- Top margin = {actual_top} EMU ({actual_top/914400:.4f} in), expected ~{EXPECTED_TOP} EMU (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 -- Top margin = {actual_top} EMU ({actual_top/914400:.4f} in), expected {EXPECTED_TOP} EMU (0.2 in)")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Bottom margin = 0.2 inches (0.25 points)
    try:
        actual_bottom = tf.margin_bottom
        if is_close(actual_bottom, EXPECTED_BOTTOM):
            print(f"PASS: Component 2 -- Bottom margin = {actual_bottom} EMU ({actual_bottom/914400:.4f} in), expected ~{EXPECTED_BOTTOM} EMU (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 -- Bottom margin = {actual_bottom} EMU ({actual_bottom/914400:.4f} in), expected {EXPECTED_BOTTOM} EMU (0.2 in)")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Left margin = 0.3 inches (0.25 points)
    try:
        actual_left = tf.margin_left
        if is_close(actual_left, EXPECTED_LEFT):
            print(f"PASS: Component 3 -- Left margin = {actual_left} EMU ({actual_left/914400:.4f} in), expected ~{EXPECTED_LEFT} EMU (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 -- Left margin = {actual_left} EMU ({actual_left/914400:.4f} in), expected {EXPECTED_LEFT} EMU (0.3 in)")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Right margin = 0.3 inches (0.25 points)
    try:
        actual_right = tf.margin_right
        if is_close(actual_right, EXPECTED_RIGHT):
            print(f"PASS: Component 4 -- Right margin = {actual_right} EMU ({actual_right/914400:.4f} in), expected ~{EXPECTED_RIGHT} EMU (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 4 -- Right margin = {actual_right} EMU ({actual_right/914400:.4f} in), expected {EXPECTED_RIGHT} EMU (0.3 in)")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
