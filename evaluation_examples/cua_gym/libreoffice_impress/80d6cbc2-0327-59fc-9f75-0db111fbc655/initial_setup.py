"""
Initial Setup: Create a 7-slide presentation with a bordered text box on slide 6
Task ID: impress_tct_085
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_085'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Text Layout Workshop"
    slide1.placeholders[1].text = "Exploring Typography and Margin Settings"

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Workshop Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Understanding text box margins and padding"
    body2.add_paragraph().text = "How internal spacing affects readability"
    body2.add_paragraph().text = "Best practices for professional layouts"
    body2.add_paragraph().text = "Hands-on exercises with LibreOffice Impress"

    # --- Slide 3: Typography Basics ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Typography Basics"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Font selection impacts visual hierarchy"
    body3.add_paragraph().text = "Line spacing controls vertical rhythm"
    body3.add_paragraph().text = "Letter spacing adjusts character density"
    body3.add_paragraph().text = "Paragraph spacing separates content blocks"

    # --- Slide 4: Margin Concepts ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Understanding Margins"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "External margins: space around the text box"
    body4.add_paragraph().text = "Internal margins: space inside the text box border"
    body4.add_paragraph().text = "Padding affects text distance from edges"
    body4.add_paragraph().text = "Consistent margins create visual harmony"

    # --- Slide 5: Layout Grid ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Layout Grid Systems"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Grid-based layouts ensure alignment across slides"
    body5.add_paragraph().text = "Column grids organize content horizontally"
    body5.add_paragraph().text = "Baseline grids maintain vertical consistency"
    body5.add_paragraph().text = "Modular grids combine both approaches"

    # --- Slide 6: Text Box with Border (KEY SLIDE) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Add a text box with visible border and default small margins
    left = Inches(2.0)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(4.5)
    txBox = slide6.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Set default small margins (python-pptx defaults are ~0.05in / 0.1in)
    # These are the default internal margins - intentionally small
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)

    # Add text content
    p1 = tf.paragraphs[0]
    p1.text = "Project Summary: Q4 Regional Performance Analysis"
    p1.alignment = PP_ALIGN.LEFT
    run1 = p1.runs[0]
    run1.font.name = "Arial"
    run1.font.size = Pt(22)
    run1.font.bold = True
    run1.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6D)

    p2 = tf.add_paragraph()
    p2.text = ("The Western region exceeded quarterly targets by 14%, driven primarily "
               "by expansion into three new metropolitan markets. Customer acquisition "
               "costs decreased by $12.40 per unit compared to Q3, while average order "
               "value increased to $287.50 across all product categories.")
    p2.alignment = PP_ALIGN.LEFT
    for run in p2.runs:
        run.font.name = "Arial"
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    p3 = tf.add_paragraph()
    p3.text = ""  # spacer paragraph

    p4 = tf.add_paragraph()
    p4.text = ("Key initiatives for Q1 include launching the digital storefront redesign, "
               "onboarding the Southeast distribution partner, and completing the CRM "
               "migration scheduled for March 15, 2026.")
    p4.alignment = PP_ALIGN.LEFT
    for run in p4.runs:
        run.font.name = "Arial"
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Add visible border to the text box
    line = txBox.line
    line.color.rgb = RGBColor(0x2E, 0x75, 0xB6)
    line.width = Pt(2.0)

    # --- Slide 7: Summary ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Key Takeaways"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "Margins improve readability and visual appeal"
    body7.add_paragraph().text = "Consistent spacing creates professional presentations"
    body7.add_paragraph().text = "Test different margin values for optimal results"
    body7.add_paragraph().text = "Apply these principles to all slide elements"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
