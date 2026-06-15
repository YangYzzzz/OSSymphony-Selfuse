"""
Initial Setup: Design a product feature showcase across slides 3-5
Task ID: impress_sales_058
Domain: libreoffice_impress

Creates an 8-slide presentation where slides 3-5 are blank.
Other slides contain realistic sales pitch content.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_058'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=None, font_name="Calibri"):
    """Helper to add a styled text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0D, 0x25, 0x3F)

    add_textbox(slide1, Inches(1.5), Inches(1.8), Inches(10), Inches(1.5),
                "NexGen Platform", font_size=44, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(2), Inches(3.5), Inches(9), Inches(1),
                "Enterprise Solutions for Tomorrow's Challenges",
                font_size=24, color=RGBColor(0xA0, 0xC4, 0xE8),
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(3.5), Inches(5.2), Inches(6), Inches(0.8),
                "Q2 2026 Product Roadmap  |  Confidential",
                font_size=14, color=RGBColor(0x80, 0x99, 0xB0),
                alignment=PP_ALIGN.CENTER)

    # --- Slide 2: Agenda / Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide2, Inches(0.8), Inches(0.5), Inches(8), Inches(1),
                "Agenda", font_size=32, bold=True,
                color=RGBColor(0x0D, 0x25, 0x3F))

    agenda_items = [
        "1. Market Overview & Competitive Landscape",
        "2. Product Feature Deep Dive",
        "3. Analytics Engine",
        "4. Security Framework",
        "5. Integration Hub",
        "6. Pricing & Licensing",
        "7. Customer Testimonials",
        "8. Next Steps & Contact",
    ]
    y_start = Inches(1.8)
    for i, item in enumerate(agenda_items):
        add_textbox(slide2, Inches(1.2), y_start + Inches(i * 0.6),
                    Inches(10), Inches(0.5), item, font_size=18,
                    color=RGBColor(0x33, 0x33, 0x33))

    # --- Slides 3, 4, 5: BLANK (feature showcase - to be designed by agent) ---
    for _ in range(3):
        prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # --- Slide 6: Pricing ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide6, Inches(0.8), Inches(0.5), Inches(8), Inches(1),
                "Pricing & Licensing", font_size=32, bold=True,
                color=RGBColor(0x0D, 0x25, 0x3F))

    tbl_shape = slide6.shapes.add_table(4, 3, Inches(1.5), Inches(2),
                                         Inches(10), Inches(3.5))
    tbl = tbl_shape.table
    headers = ["Plan", "Monthly Price", "Included Users"]
    data_rows = [
        ["Starter", "$499/mo", "Up to 25"],
        ["Professional", "$1,299/mo", "Up to 100"],
        ["Enterprise", "Custom", "Unlimited"],
    ]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x0D, 0x25, 0x3F)
    for r, row in enumerate(data_rows, 1):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(14)

    # --- Slide 7: Testimonials ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide7, Inches(0.8), Inches(0.5), Inches(8), Inches(1),
                "What Our Customers Say", font_size=32, bold=True,
                color=RGBColor(0x0D, 0x25, 0x3F))

    testimonials = [
        ('"NexGen\'s analytics cut our reporting time by 60%. The dashboards are intuitive and the real-time data feeds are incredibly reliable."',
         "— Rachel Torres, VP of Operations, Meridian Corp"),
        ('"We migrated from three separate tools to NexGen Integration Hub. The unified workflow saved us $340K annually."',
         "— David Okafor, CTO, Pinnacle Health Systems"),
        ('"The security framework gave us SOC 2 compliance in half the expected timeline. Outstanding support team."',
         "— Anya Petrova, CISO, Lumina Financial"),
    ]
    for i, (quote, attribution) in enumerate(testimonials):
        y = Inches(1.8) + Inches(i * 1.8)
        add_textbox(slide7, Inches(1.2), y, Inches(10.5), Inches(1),
                    quote, font_size=16, color=RGBColor(0x33, 0x33, 0x33))
        add_textbox(slide7, Inches(1.2), y + Inches(0.9), Inches(10), Inches(0.5),
                    attribution, font_size=13, bold=True,
                    color=RGBColor(0x0D, 0x25, 0x3F))

    # --- Slide 8: Contact / Closing ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide8.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0D, 0x25, 0x3F)

    add_textbox(slide8, Inches(2), Inches(2), Inches(9), Inches(1.2),
                "Ready to Get Started?", font_size=40, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_textbox(slide8, Inches(2.5), Inches(3.5), Inches(8), Inches(0.8),
                "Contact our enterprise sales team for a personalized demo.",
                font_size=20, color=RGBColor(0xA0, 0xC4, 0xE8),
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide8, Inches(3), Inches(4.8), Inches(7), Inches(0.5),
                "sales@nexgenplatform.com  |  +1 (415) 555-0192",
                font_size=16, color=RGBColor(0x80, 0x99, 0xB0),
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide8, Inches(3), Inches(5.5), Inches(7), Inches(0.5),
                "www.nexgenplatform.com",
                font_size=16, color=RGBColor(0x80, 0x99, 0xB0),
                alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
