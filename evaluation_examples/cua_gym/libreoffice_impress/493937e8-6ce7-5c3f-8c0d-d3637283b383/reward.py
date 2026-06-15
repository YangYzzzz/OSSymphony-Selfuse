"""
Reward Script: Product feature showcase on slides 3-5 with circle icons, feature names, descriptions, bullets, and animations.
Task ID: impress_sales_058
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Circle shapes with correct letters (A/S/I) on slides 3-5
  Component 2 (0.30): Feature name text in ~28pt bold on slides 3-5
  Component 3 (0.20): Description text (~16pt) with 3+ sentences on slides 3-5
  Component 4 (0.20): Entrance animations present on slides 3-5
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_058'
FILE_PATH = os.path.join(WORKDIR, f'{TASK_ID}.pptx')

# Expected feature data per slide (0-indexed: slide index 2,3,4 = slides 3,4,5)
EXPECTED_SLIDES = {
    2: {'letter': 'A', 'feature': 'Analytics'},
    3: {'letter': 'S', 'feature': 'Security'},
    4: {'letter': 'I', 'feature': 'Integration'},
}


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Circle shapes with correct letters A/S/I (0.30 points, 0.10 per slide)
    for slide_idx, expected in EXPECTED_SLIDES.items():
        try:
            slide = prs.slides[slide_idx]
            letter_found = False
            for shape in slide.shapes:
                # Check for auto shape (oval/circle)
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    if shape.has_text_frame:
                        shape_text = shape.text_frame.text.strip()
                        if shape_text == expected['letter']:
                            # Verify it has a solid fill (colored circle)
                            try:
                                fill_type = shape.fill.type
                                if fill_type is not None:
                                    letter_found = True
                                else:
                                    letter_found = True  # shape exists with letter, fill may be theme-based
                            except Exception:
                                letter_found = True  # shape exists with correct letter
                            break
            if letter_found:
                print(f"PASS: Component 1 - Slide {slide_idx+1} has circle with letter '{expected['letter']}' (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 1 - Slide {slide_idx+1} missing circle shape with letter '{expected['letter']}'")
        except Exception as e:
            print(f"ERROR: Component 1 - Slide {slide_idx+1}: {e}")

    # Component 2: Feature name in ~28pt bold (0.30 points, 0.10 per slide)
    for slide_idx, expected in EXPECTED_SLIDES.items():
        try:
            slide = prs.slides[slide_idx]
            feature_found = False
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        # Check if the text contains the feature name (case-insensitive)
                        if expected['feature'].lower() in text.lower() and len(text) < 50:
                            # Verify bold and approximate font size (~28pt = 355600 EMU, tolerance +/- 20%)
                            for run in para.runs:
                                if run.text.strip():
                                    is_bold = run.font.bold is True
                                    font_size = run.font.size
                                    # 28pt = 355600 EMU; allow 20pt-36pt range (254000-457200)
                                    size_ok = (font_size is not None and 254000 <= font_size <= 457200)
                                    if is_bold and size_ok:
                                        feature_found = True
                                        print(f"PASS: Component 2 - Slide {slide_idx+1} has feature name '{text}' in {font_size/12700:.0f}pt bold (0.10 pts)")
                                        break
                        if feature_found:
                            break
                if feature_found:
                    break
            if feature_found:
                total_score += 0.10
            else:
                print(f"FAIL: Component 2 - Slide {slide_idx+1} missing feature name '{expected['feature']}' in ~28pt bold")
        except Exception as e:
            print(f"ERROR: Component 2 - Slide {slide_idx+1}: {e}")

    # Component 3: Description text in ~16pt with multi-sentence content AND 3 bullet sub-features (0.20 points, ~0.067 per slide)
    per_slide_c3 = round(0.20 / 3, 4)
    for slide_idx, expected in EXPECTED_SLIDES.items():
        try:
            slide = prs.slides[slide_idx]
            desc_found = False
            bullets_found = False
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shape.has_text_frame:
                    paras = [p for p in shape.text_frame.paragraphs if p.text.strip()]
                    # Check for description: single paragraph with longer text (~16pt)
                    if len(paras) == 1 and len(paras[0].text.strip()) > 80:
                        for run in paras[0].runs:
                            if run.text.strip() and run.font.size is not None:
                                # 16pt = 203200 EMU; allow 12pt-20pt (152400-254000)
                                if 152400 <= run.font.size <= 254000:
                                    desc_found = True
                                    break
                    # Check for bullets: 3 paragraphs with text
                    if len(paras) >= 3 and all(len(p.text.strip()) > 10 for p in paras[:3]):
                        # Ensure these aren't the description (each bullet should be a separate paragraph)
                        max_text_len = max(len(p.text.strip()) for p in paras[:3])
                        if max_text_len < 200:  # bullets are shorter than descriptions
                            bullets_found = True

            if desc_found and bullets_found:
                print(f"PASS: Component 3 - Slide {slide_idx+1} has description and 3 bullet sub-features ({per_slide_c3} pts)")
                total_score += per_slide_c3
            elif desc_found:
                print(f"PARTIAL: Component 3 - Slide {slide_idx+1} has description but missing 3 bullets ({per_slide_c3/2:.4f} pts)")
                total_score += per_slide_c3 / 2
            elif bullets_found:
                print(f"PARTIAL: Component 3 - Slide {slide_idx+1} has bullets but missing description ({per_slide_c3/2:.4f} pts)")
                total_score += per_slide_c3 / 2
            else:
                print(f"FAIL: Component 3 - Slide {slide_idx+1} missing description and/or bullet sub-features")
        except Exception as e:
            print(f"ERROR: Component 3 - Slide {slide_idx+1}: {e}")

    # Component 4: Entrance animations on slides 3-5 (0.20 points, ~0.067 per slide)
    per_slide_c4 = round(0.20 / 3, 4)
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            for slide_idx in [2, 3, 4]:
                slide_num = slide_idx + 1
                fname = f'ppt/slides/slide{slide_num}.xml'
                try:
                    with zf.open(fname) as f:
                        root = ET.parse(f).getroot()
                        timing = root.find('.//p:timing', ns)
                        if timing is not None:
                            # Check for entrance animation elements (presetClass="entr")
                            timing_xml = ET.tostring(timing, encoding='unicode')
                            if 'presetClass="entr"' in timing_xml:
                                print(f"PASS: Component 4 - Slide {slide_num} has entrance animations ({per_slide_c4} pts)")
                                total_score += per_slide_c4
                            else:
                                print(f"FAIL: Component 4 - Slide {slide_num} has timing but no entrance animations")
                        else:
                            print(f"FAIL: Component 4 - Slide {slide_num} has no animation timing")
                except KeyError:
                    print(f"FAIL: Component 4 - Slide {slide_num} XML not found")
    except Exception as e:
        print(f"ERROR: Component 4 - Could not read ZIP for animations: {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.4f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(FILE_PATH)
