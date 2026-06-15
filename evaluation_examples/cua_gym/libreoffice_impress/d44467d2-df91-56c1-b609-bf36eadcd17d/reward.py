"""
Reward Script: Strikethrough first and second bullet points on slide 4
Task ID: osworld_impress_strikethrough_text_003
Domain: libreoffice_impress
Scoring:
  Component 1: Bullet 1 (para 0) has sngStrike             — 0.5 points
  Component 2: Bullet 2 (para 1) has sngStrike             — 0.3 points
  Component 3: Bullets 3-5 remain noStrike AND bullets 1-2 are struck (compound integrity) — 0.2 points
Total: 1.0

Note: Component 3 is a compound check requiring BOTH bullets 1&2 struck AND bullets 3-5 not struck.
It only passes on golden (where 1&2 have strikethrough). On initial all bullets are noStrike, so
the compound condition fails because bullets 1&2 don't have strikethrough.
"""

import os

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_strikethrough_text_003'


def persist_app_state():
    """Send Ctrl+S to persist any unsaved LibreOffice Impress edits."""
    try:
        os.environ["DISPLAY"] = ":0"
        import pyautogui
        import time
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def get_strike_for_para(para):
    """Return the strike attribute for the first non-empty run in a paragraph.
    Falls back to checking via XML if no runs found.
    Returns 'noStrike', 'sngStrike', or 'dblStrike'.
    """
    nonempty = [r for r in para.runs if (r.text or "").strip()]
    if nonempty:
        return nonempty[0].font._element.attrib.get('strike', 'noStrike')
    # Fallback: check all runs including empty ones
    for run in para.runs:
        val = run.font._element.attrib.get('strike', None)
        if val is not None:
            return val
    return 'noStrike'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0

    Task: Strikethrough the first and second bullet points on slide 4.
    Slide 4 (index 3) contains a 5-bullet to-do list in Content Placeholder 2.
    Only bullets 1 and 2 (paragraph indices 0 and 1) should have sngStrike applied.
    Bullets 3-5 must remain without strikethrough.
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"CRITICAL: Expected at least 4 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide4 = prs.slides[3]  # Slide 4 is index 3

    # Locate the content placeholder with the bullet list
    content_shape = None
    for shape in slide4.shapes:
        if shape.has_text_frame and shape.name == "Content Placeholder 2":
            content_shape = shape
            break

    if content_shape is None:
        # Fallback: look for a shape with 5+ paragraphs
        for shape in slide4.shapes:
            if shape.has_text_frame and len(shape.text_frame.paragraphs) >= 5:
                content_shape = shape
                break

    if content_shape is None:
        print("CRITICAL: Could not find content shape with bullet list on slide 4")
        print("REWARD: 0.0")
        return 0.0

    paras = content_shape.text_frame.paragraphs

    # Collect strike values for all 5 bullet paragraphs
    strike_values = []
    for i in range(min(5, len(paras))):
        strike_values.append(get_strike_for_para(paras[i]))

    # Component 1: Bullet 1 (paragraph index 0) has strikethrough (0.5 points)
    # This FAILS on initial (noStrike) and PASSES on golden (sngStrike)
    try:
        if len(strike_values) >= 1:
            sval = strike_values[0]
            bullet1_text = paras[0].text.strip()
            if sval in ('sngStrike', 'dblStrike'):
                print(f"PASS: Component 1 — Bullet 1 has strikethrough (strike={sval}), "
                      f"text='{bullet1_text}' (0.5 pts)")
                total_score += 0.5
            else:
                print(f"FAIL: Component 1 — Bullet 1 expected strikethrough, got strike={sval}, "
                      f"text='{bullet1_text}'")
        else:
            print("FAIL: Component 1 — No paragraphs found in content shape")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Bullet 2 (paragraph index 1) has strikethrough (0.3 points)
    # This FAILS on initial (noStrike) and PASSES on golden (sngStrike)
    try:
        if len(strike_values) >= 2:
            sval = strike_values[1]
            bullet2_text = paras[1].text.strip()
            if sval in ('sngStrike', 'dblStrike'):
                print(f"PASS: Component 2 — Bullet 2 has strikethrough (strike={sval}), "
                      f"text='{bullet2_text}' (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 2 — Bullet 2 expected strikethrough, got strike={sval}, "
                      f"text='{bullet2_text}'")
        else:
            print("FAIL: Component 2 — Less than 2 paragraphs found in content shape")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Compound check — bullets 1&2 are struck THROUGH AND bullets 3-5 remain noStrike (0.2 points)
    # This is a compound integrity check: passes ONLY when BOTH conditions are true together.
    # On initial: bullets 1&2 have noStrike, so the sub-condition "bullets 1&2 struck" fails → component fails → 0 pts
    # On golden: bullets 1&2 have sngStrike AND bullets 3-5 have noStrike → both sub-conditions pass → 0.2 pts
    try:
        if len(strike_values) >= 5:
            bullets_12_struck = all(strike_values[i] in ('sngStrike', 'dblStrike') for i in range(2))
            bullets_35_clean = all(strike_values[i] not in ('sngStrike', 'dblStrike') for i in range(2, 5))
            if bullets_12_struck and bullets_35_clean:
                print("PASS: Component 3 — Bullets 1-2 struck through AND bullets 3-5 remain clean (0.2 pts)")
                total_score += 0.2
            else:
                if not bullets_12_struck:
                    print("FAIL: Component 3 — Compound check: bullets 1-2 are not all struck through")
                if not bullets_35_clean:
                    print("FAIL: Component 3 — Compound check: one or more of bullets 3-5 unexpectedly struck through")
        else:
            print(f"FAIL: Component 3 — Expected 5 bullet paragraphs, found {len(strike_values)}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
persist_app_state()

if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
