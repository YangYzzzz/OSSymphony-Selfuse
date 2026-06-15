"""
Initial Setup: 5-slide team meeting presentation with a to-do list on slide 4 (no strikethrough).
Task ID: osworld_impress_strikethrough_text_003
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_strikethrough_text_003'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q2 Team Meeting"
    slide1.placeholders[1].text = "April 2025 | Engineering Division"

    # --- Slide 2: Team Updates ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Team Updates"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Sarah Chen promoted to Senior Engineer"
    updates = [
        "Marcus Johnson joined the infrastructure team",
        "Priya Patel completed AWS certification",
        "DevOps migration is 80% complete",
        "Q1 performance reviews scheduled for April 18",
    ]
    for update in updates:
        p = tf2.add_paragraph()
        p.text = update
        p.level = 0

    # --- Slide 3: Agenda ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Meeting Agenda"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "1. Team Updates"
    agenda_items = [
        "2. Q2 Planning & Roadmap",
        "3. Action Items Review",
        "4. Budget Discussion",
        "5. Open Floor / Q&A",
    ]
    for item in agenda_items:
        p = tf3.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 4: Action Items (To-Do List) — NO strikethrough ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Action Items"
    tf4 = slide4.placeholders[1].text_frame
    tf4.word_wrap = True

    todo_items = [
        "Finalize Q2 roadmap document and share with stakeholders",
        "Schedule 1:1 check-ins with all direct reports",
        "Submit budget proposal to finance by April 25",
        "Update project tracker with current milestone status",
        "Coordinate with design team on new UI mockups",
    ]

    tf4.text = todo_items[0]
    run0 = tf4.paragraphs[0].runs[0]
    run0.font.size = Pt(20)
    # No strikethrough on any bullets in initial state

    for item in todo_items[1:]:
        p = tf4.add_paragraph()
        p.text = item
        p.level = 0
        run = p.runs[0]
        run.font.size = Pt(20)
        # No strikethrough

    # --- Slide 5: Wrap-Up ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Wrap-Up & Next Steps"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Next meeting: May 7, 2025"
    wrap_items = [
        "All action items due by April 30",
        "Weekly stand-ups every Tuesday at 10am",
        "Questions? Reach out on Slack #team-eng",
    ]
    for item in wrap_items:
        p = tf5.add_paragraph()
        p.text = item
        p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
