"""
FINAL REWARD SCRIPT - SUCCESS
Task: Every time I drag our logo onto slide 47 it shows up huge or squished. In LibreOffice Impress, how do I insert the image at ~/Desktop/logo.png so it lands at exactly 4 cm wide and the aspect ratio stays locked?
Generated: 2025-09-10 23:16:05
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
import io
from pptx import Presentation
from PIL import Image


def verify_logo_slide(
    file_path: str,
    slide_index: int = 46,  # zero-based → slide 47
    target_width_cm: float = 4.0,
    width_tol_cm: float = 0.05,  # ±0.05 cm tolerance (≈0.5 mm)
    ratio_tol: float = 0.01      # ≤1 % distortion allowed
) -> float:
    """Verify that an image on a specific slide is exactly the required width
    (within tolerance) and that its aspect ratio is preserved.

    Returns a progressive score between 0.0 and 1.0:
    • 0.5 points for correct width
    • 0.5 points for preserved aspect ratio
    """

    # Convert centimetres to EMUs (English Metric Units)
    CM_TO_EMU = 360_000  # 1 cm = 360 000 EMU
    target_width_emu = target_width_cm * CM_TO_EMU
    width_tol_emu = width_tol_cm * CM_TO_EMU

    # 1) Load the presentation ------------------------------------------------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Unable to open presentation: {e}")
        return 0.0
    print(f"✓ Presentation loaded – {len(prs.slides)} slides found")

    # 2) Check slide index -----------------------------------------------------
    if slide_index >= len(prs.slides):
        print(
            f"✗ Slide {slide_index + 1} does not exist (presentation has {len(prs.slides)} slides)"
        )
        return 0.0
    slide = prs.slides[slide_index]

    # 3) Locate picture shapes -------------------------------------------------
    picture_shapes = [sh for sh in slide.shapes if sh.shape_type == 13]  # PICTURE
    print(f"Found {len(picture_shapes)} picture shape(s) on slide {slide_index + 1}")

    if not picture_shapes:
        print("✗ No images found on the target slide")
        return 0.0

    # 4) Evaluate each picture -------------------------------------------------
    best_width_score = 0.0
    best_ratio_score = 0.0

    for idx, pic in enumerate(picture_shapes, start=1):
        w_emu, h_emu = pic.width, pic.height
        print(f" Picture {idx}: width={w_emu} EMU, height={h_emu} EMU")

        # Width verification --------------------------------------------------
        if abs(w_emu - target_width_emu) <= width_tol_emu:
            best_width_score = 0.5  # award width points once at most
            print("  ✓ Width within tolerance (0.5 points)")
        else:
            print("  ✗ Width outside tolerance")

        # Aspect-ratio verification ------------------------------------------
        try:
            img = Image.open(io.BytesIO(pic.image.blob))
            orig_ratio = img.width / img.height if img.height else None
            slide_ratio = w_emu / h_emu if h_emu else None
            if (
                orig_ratio
                and slide_ratio
                and abs(slide_ratio - orig_ratio) / orig_ratio <= ratio_tol
            ):
                best_ratio_score = 0.5  # award ratio points once at most
                print("  ✓ Aspect ratio preserved (0.5 points)")
            else:
                print("  ✗ Aspect ratio distorted")
        except Exception as e:
            print(f"  ✗ Could not analyse embedded image: {e}")

    # 5) Final scoring ---------------------------------------------------------
    total_score = best_width_score + best_ratio_score
    final_score = min(total_score, 1.0)
    print(f"Total score: {final_score}")
    return final_score


if __name__ == "__main__":
    # Path to the candidate presentation (provided by the grading environment)
    candidate_path = (
        "/home/user/"
        "every_time_i_drag_our_logo_onto_slide_47_it_shows_up_huge_or_"
        "squished_in_libreoffice_impress_how_do__golden.pptx"
    )

    reward = verify_logo_slide(candidate_path)
    print(f"REWARD: {reward}")

