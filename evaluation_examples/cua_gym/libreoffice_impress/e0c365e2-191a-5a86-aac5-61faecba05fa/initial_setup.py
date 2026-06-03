"""
Initial Setup: Create a Development Lifecycle presentation with 8 slides.
Slide 3 has title 'Our Process' but is otherwise empty (no chevrons).
Task ID: impress_rp_021
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_021'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_text(slide, text):
    """Set the title placeholder text on a slide."""
    if slide.shapes.title:
        slide.shapes.title.text = text


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Add a textbox with formatted text to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Development Lifecycle"
    slide1.placeholders[1].text = "Streamlining Software Delivery at TechForward Inc."

    # --- Slide 2: Project Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Project Overview"
    body2 = slide2.placeholders[1]
    tf2 = body2.text_frame
    tf2.text = "TechForward Inc. has adopted an agile development methodology to accelerate product releases while maintaining quality standards."
    p2 = tf2.add_paragraph()
    p2.text = "Our team of 45 engineers across 6 squads delivers bi-weekly sprints with a 97.3% on-time delivery rate."
    p3 = tf2.add_paragraph()
    p3.text = "This presentation outlines our end-to-end development lifecycle and key performance metrics."

    # --- Slide 3: Our Process (EMPTY body - task target) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title as a textbox since blank layout has no title placeholder
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Our Process"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x52, 0x76)
    # No other shapes - this is what the agent needs to fill

    # --- Slide 4: Team Structure ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Team Structure"
    body4 = slide4.placeholders[1]
    tf4 = body4.text_frame
    tf4.text = "Frontend Squad: 8 engineers led by Priya Sharma"
    for line in [
        "Backend Squad: 10 engineers led by David Kim",
        "DevOps Squad: 6 engineers led by Maria Santos",
        "QA Squad: 7 engineers led by James Wright",
        "Data Squad: 8 engineers led by Aisha Patel",
        "Security Squad: 6 engineers led by Tom Nakamura",
    ]:
        p = tf4.add_paragraph()
        p.text = line

    # --- Slide 5: Technology Stack ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Technology Stack"
    body5 = slide5.placeholders[1]
    tf5 = body5.text_frame
    tf5.text = "Frontend: React 18, TypeScript, Tailwind CSS"
    for line in [
        "Backend: Python 3.12, FastAPI, PostgreSQL 16",
        "Infrastructure: AWS EKS, Terraform, ArgoCD",
        "Monitoring: Datadog, PagerDuty, Grafana",
        "CI/CD: GitHub Actions, Docker, Helm Charts",
    ]:
        p = tf5.add_paragraph()
        p.text = line

    # --- Slide 6: Sprint Metrics ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Sprint Metrics Q1 2025"
    body6 = slide6.placeholders[1]
    tf6 = body6.text_frame
    tf6.text = "Average velocity: 142 story points per sprint"
    for line in [
        "Defect escape rate: 0.8% (target < 2%)",
        "Code coverage: 94.2% across all repositories",
        "Mean time to recovery: 12 minutes",
        "Deployment frequency: 8.3 deploys per day",
    ]:
        p = tf6.add_paragraph()
        p.text = line

    # --- Slide 7: Roadmap ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "2025 Product Roadmap"
    body7 = slide7.placeholders[1]
    tf7 = body7.text_frame
    tf7.text = "Q1: Platform migration to microservices architecture"
    for line in [
        "Q2: Launch AI-powered recommendation engine",
        "Q3: International expansion (EU, APAC markets)",
        "Q4: Real-time collaboration features release",
    ]:
        p = tf7.add_paragraph()
        p.text = line

    # --- Slide 8: Contact & Next Steps ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Next Steps & Contact"
    body8 = slide8.placeholders[1]
    tf8 = body8.text_frame
    tf8.text = "Review sprint retrospective findings by March 28, 2025"
    for line in [
        "Schedule architecture review with platform team",
        "Finalize Q2 hiring plan for 12 new positions",
        "Contact: engineering@techforward.com",
    ]:
        p = tf8.add_paragraph()
        p.text = line

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
