"""
Reward Script: Verify chevron process flow on slide 3
Task ID: impress_rp_021
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Five chevron/arrow shapes exist on slide 3
  Component 2 (0.25): Correct labels in left-to-right order
  Component 3 (0.25): Correct fill color gradient from light to dark blue
  Component 4 (0.25): White, bold, 14pt text formatting on all chevrons
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_021'

# Expected values from task description and context
EXPECTED_LABELS = ['Research', 'Design', 'Develop', 'Test', 'Deploy']
EXPECTED_COLORS = ['AED6F1', '7FB3D8', '5499C7', '2E86C1', '1A5276']
EXPECTED_TEXT_COLOR = 'FFFFFF'
EXPECTED_FONT_SIZE = Pt(14)  # 177800 EMU


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]

    # Collect chevron/auto shapes from slide 3 (exclude placeholders and text boxes)
    chevron_shapes = []
    for shape in slide3.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            chevron_shapes.append(shape)

    # Sort by left position (left to right ordering)
    chevron_shapes.sort(key=lambda s: s.left)

    # Component 1: Five chevron/arrow shapes exist on slide 3 (0.25 points)
    try:
        num_chevrons = len(chevron_shapes)
        if num_chevrons == 5:
            print(f"PASS: Component 1 — Found exactly 5 auto shapes on slide 3 (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — Expected 5 auto shapes on slide 3, found {num_chevrons}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # If we don't have exactly 5 chevrons, remaining checks use what we have
    if len(chevron_shapes) < 5:
        print(f"Cannot fully verify remaining components with only {len(chevron_shapes)} shapes")
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: Correct labels in left-to-right order (0.25 points)
    try:
        actual_labels = []
        for shape in chevron_shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                actual_labels.append(text)
            else:
                actual_labels.append('')

        matching_labels = sum(
            1 for a, e in zip(actual_labels, EXPECTED_LABELS)
            if a.lower() == e.lower()
        )
        if matching_labels == 5:
            print(f"PASS: Component 2 — All 5 labels correct in order: {actual_labels} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 — Labels mismatch. Expected {EXPECTED_LABELS}, found {actual_labels} ({matching_labels}/5 correct)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Correct fill color gradient (0.25 points)
    try:
        actual_colors = []
        color_matches = 0
        for i, shape in enumerate(chevron_shapes):
            fill = shape.fill
            if fill.type is not None and fill.type == 1:  # SOLID fill
                try:
                    rgb = str(fill.fore_color.rgb).upper()
                    actual_colors.append(rgb)
                    if rgb == EXPECTED_COLORS[i].upper():
                        color_matches += 1
                except Exception:
                    actual_colors.append('ERROR')
            else:
                actual_colors.append(f'type={fill.type}')

        if color_matches == 5:
            print(f"PASS: Component 3 — All 5 fill colors correct: {actual_colors} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 — Fill colors mismatch. Expected {EXPECTED_COLORS}, found {actual_colors} ({color_matches}/5 correct)")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: White, bold, 14pt text formatting (0.25 points)
    try:
        formatting_ok = 0
        formatting_details = []
        for i, shape in enumerate(chevron_shapes):
            if not shape.has_text_frame:
                formatting_details.append(f"Shape {i}: no text frame")
                continue
            runs = [r for p in shape.text_frame.paragraphs for r in p.runs if (r.text or '').strip()]
            if not runs:
                formatting_details.append(f"Shape {i}: no runs")
                continue

            shape_ok = True
            for run in runs:
                # Check bold (True required; None means inherited = not explicitly bold)
                is_bold = run.font.bold is True
                if not is_bold:
                    shape_ok = False
                    formatting_details.append(f"Shape {i} '{run.text}': bold={run.font.bold}, expected True")

                # Check font size (14pt = 177800 EMU)
                size_ok = False
                if run.font.size is not None:
                    # Allow small tolerance for size
                    if abs(run.font.size - EXPECTED_FONT_SIZE) <= 1000:
                        size_ok = True
                if not size_ok:
                    shape_ok = False
                    formatting_details.append(f"Shape {i} '{run.text}': size={run.font.size}, expected {EXPECTED_FONT_SIZE}")

                # Check white text color
                try:
                    if run.font.color.type is not None:
                        rgb = str(run.font.color.rgb).upper()
                        if rgb != EXPECTED_TEXT_COLOR:
                            shape_ok = False
                            formatting_details.append(f"Shape {i} '{run.text}': color={rgb}, expected {EXPECTED_TEXT_COLOR}")
                    else:
                        shape_ok = False
                        formatting_details.append(f"Shape {i} '{run.text}': color type is None")
                except Exception as ce:
                    shape_ok = False
                    formatting_details.append(f"Shape {i} '{run.text}': color error: {ce}")

            if shape_ok:
                formatting_ok += 1

        if formatting_ok == 5:
            print(f"PASS: Component 4 — All 5 shapes have white, bold, 14pt text (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 4 — Formatting issues in {5 - formatting_ok} shapes:")
            for detail in formatting_details:
                print(f"  {detail}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
