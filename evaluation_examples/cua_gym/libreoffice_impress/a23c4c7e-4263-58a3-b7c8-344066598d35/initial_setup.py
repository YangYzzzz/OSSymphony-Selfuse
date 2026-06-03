"""
Initial Setup: Create presentation with column chart on slide 2 (auto-scaled Y-axis)
Task ID: impress_tct_058
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_058'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Quarterly Sales Performance"
    slide1.placeholders[1].text = "FY 2025 Regional Breakdown"

    # --- Slide 2: Chart Slide ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Add a title text box
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Regional Sales by Quarter (in thousands)"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x6E)

    # Chart data — values range roughly 60–380 so auto-scale lands around 50–400
    chart_data = CategoryChartData()
    chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
    chart_data.add_series('North', (125, 210, 185, 240))
    chart_data.add_series('South', (95, 145, 170, 200))
    chart_data.add_series('East', (180, 260, 310, 375))
    chart_data.add_series('West', (65, 110, 150, 195))

    chart_shape = slide2.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.5),
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = True

    # --- Slide 3: Key Insights ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Key Insights"
    p3.alignment = PP_ALIGN.LEFT
    r3 = p3.runs[0]
    r3.font.size = Pt(28)
    r3.font.bold = True
    r3.font.color.rgb = RGBColor(0x2E, 0x4A, 0x6E)

    body3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.0))
    btf = body3.text_frame
    btf.word_wrap = True
    insights = [
        "East region showed the strongest growth trajectory, increasing 108% from Q1 to Q4.",
        "All regions demonstrated positive quarter-over-quarter growth in Q4.",
        "South region maintained steady but moderate growth throughout the year.",
        "West region recovered strongly in Q3 and Q4 after a slow Q1 start.",
        "Total company revenue across all regions exceeded $2.9 million for the fiscal year.",
    ]
    for i, text in enumerate(insights):
        if i == 0:
            p = btf.paragraphs[0]
        else:
            p = btf.add_paragraph()
        p.text = text
        p.space_after = Pt(12)
        for run in p.runs:
            run.font.size = Pt(16)

    # --- Slide 4: Next Steps ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Next Steps"
    p4.alignment = PP_ALIGN.LEFT
    r4 = p4.runs[0]
    r4.font.size = Pt(28)
    r4.font.bold = True
    r4.font.color.rgb = RGBColor(0x2E, 0x4A, 0x6E)

    body4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(4.0))
    btf4 = body4.text_frame
    btf4.word_wrap = True
    steps = [
        "Expand East region sales team by 15% to capitalize on growth momentum.",
        "Launch targeted marketing campaign for South region in Q1 2026.",
        "Conduct West region customer satisfaction survey to sustain Q3–Q4 gains.",
        "Set regional targets aligned with overall company growth objective of 25%.",
    ]
    for i, text in enumerate(steps):
        if i == 0:
            p = btf4.paragraphs[0]
        else:
            p = btf4.add_paragraph()
        p.text = text
        p.space_after = Pt(12)
        for run in p.runs:
            run.font.size = Pt(16)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
