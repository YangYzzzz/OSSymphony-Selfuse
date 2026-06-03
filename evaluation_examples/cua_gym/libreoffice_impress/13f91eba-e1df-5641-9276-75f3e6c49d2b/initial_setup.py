"""
Initial Setup: Animated countdown sequence on slide 1
Task ID: impress_gf2_015
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_015'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_centered_textbox(slide, text, top, font_size, font_color, bold=True,
                         width=Inches(3), height=Inches(1.5),
                         slide_width=Inches(10)):
    """Add a centered text box with specified properties."""
    left = (slide_width - width) // 2
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = font_color
    return txBox


def create_initial():
    prs = Presentation()
    # Standard widescreen
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========== Slide 1: Countdown + Title ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only layout

    # Text box '3' - centered, upper area
    add_centered_textbox(slide1, '3', Inches(2.0), Pt(120),
                         RGBColor(0xE7, 0x4C, 0x3C), bold=True,
                         width=Inches(2), height=Inches(1.8))

    # Text box '2' - centered, same position (stacked)
    add_centered_textbox(slide1, '2', Inches(2.0), Pt(120),
                         RGBColor(0xF3, 0x9C, 0x12), bold=True,
                         width=Inches(2), height=Inches(1.8))

    # Text box '1' - centered, same position (stacked)
    add_centered_textbox(slide1, '1', Inches(2.0), Pt(120),
                         RGBColor(0x27, 0xAE, 0x60), bold=True,
                         width=Inches(2), height=Inches(1.8))

    # Main title at bottom
    add_centered_textbox(slide1, 'Annual Gala 2024', Inches(5.5), Pt(44),
                         RGBColor(0x2C, 0x3E, 0x50), bold=True,
                         width=Inches(8), height=Inches(1.2))

    # Set dark background for slide 1
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # ========== Slide 2: Welcome ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    add_centered_textbox(slide2, 'Welcome to the Annual Gala', Inches(1.0),
                         Pt(36), RGBColor(0xFF, 0xFF, 0xFF), bold=True,
                         width=Inches(8), height=Inches(1.0))
    add_centered_textbox(slide2, 'December 14, 2024 | Grand Ballroom | 7:00 PM',
                         Inches(2.5), Pt(20), RGBColor(0xBD, 0xC3, 0xC7),
                         bold=False, width=Inches(8), height=Inches(0.8))

    # ========== Slide 3: Agenda ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_centered_textbox(slide3, 'Evening Agenda', Inches(0.5), Pt(32),
                         RGBColor(0x2C, 0x3E, 0x50), bold=True,
                         width=Inches(8), height=Inches(1.0))
    agenda_box = slide3.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(4))
    tf = agenda_box.text_frame
    tf.word_wrap = True
    items = [
        '7:00 PM - Cocktail Reception & Networking',
        '7:45 PM - Opening Remarks by CEO Patricia Harmon',
        '8:00 PM - Awards Ceremony',
        '8:45 PM - Keynote Speaker: Dr. Amara Osei',
        '9:15 PM - Dinner Service',
        '10:00 PM - Live Music & Dancing',
    ]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)
        p.space_after = Pt(12)

    # ========== Slide 4: Award Categories ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_centered_textbox(slide4, 'Award Categories', Inches(0.5), Pt(32),
                         RGBColor(0xE7, 0x4C, 0x3C), bold=True,
                         width=Inches(8), height=Inches(1.0))
    awards_box = slide4.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(7), Inches(4))
    tf = awards_box.text_frame
    tf.word_wrap = True
    awards = [
        'Employee of the Year',
        'Innovation Champion',
        'Customer Excellence Award',
        'Rising Star Award',
        'Team Collaboration Award',
    ]
    for i, award in enumerate(awards):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f'  {award}'
        run.font.size = Pt(22)
        run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
        p.space_after = Pt(14)

    # ========== Slide 5: Keynote Speaker ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    fill5 = slide5.background.fill
    fill5.solid()
    fill5.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    add_centered_textbox(slide5, 'Keynote Speaker', Inches(0.5), Pt(28),
                         RGBColor(0xF3, 0x9C, 0x12), bold=True,
                         width=Inches(8), height=Inches(0.8))
    add_centered_textbox(slide5, 'Dr. Amara Osei', Inches(2.0), Pt(36),
                         RGBColor(0xFF, 0xFF, 0xFF), bold=True,
                         width=Inches(8), height=Inches(1.0))
    add_centered_textbox(slide5, '"The Future of Sustainable Innovation"\n\n'
                         'Dr. Osei is a renowned expert in sustainable technology\n'
                         'and has led groundbreaking research at Stanford University.',
                         Inches(3.5), Pt(16), RGBColor(0xBD, 0xC3, 0xC7),
                         bold=False, width=Inches(7), height=Inches(2.5))

    # ========== Slide 6: Past Highlights ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_centered_textbox(slide6, 'Highlights from Last Year', Inches(0.5),
                         Pt(32), RGBColor(0x2C, 0x3E, 0x50), bold=True,
                         width=Inches(8), height=Inches(1.0))
    stats_data = [
        ('$2.4M', 'Raised for Charity'),
        ('350+', 'Attendees'),
        ('12', 'Awards Presented'),
    ]
    for i, (num, label) in enumerate(stats_data):
        left = Inches(1.0 + i * 3.0)
        add_centered_textbox(slide6, num, Inches(2.5), Pt(48),
                             RGBColor(0xE7, 0x4C, 0x3C), bold=True,
                             width=Inches(2.5), height=Inches(1.2))
        # Reposition the last added shape
        prs.slides[5].shapes[-1].left = left
        add_centered_textbox(slide6, label, Inches(3.8), Pt(16),
                             RGBColor(0x7F, 0x8C, 0x8D), bold=False,
                             width=Inches(2.5), height=Inches(0.6))
        prs.slides[5].shapes[-1].left = left

    # ========== Slide 7: Sponsors ==========
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_centered_textbox(slide7, 'Our Generous Sponsors', Inches(0.5),
                         Pt(32), RGBColor(0x2C, 0x3E, 0x50), bold=True,
                         width=Inches(8), height=Inches(1.0))
    sponsors = ['Meridian Technologies', 'Atlas Financial Group',
                'Cascade Health Systems', 'Pinnacle Media Corp']
    sponsor_box = slide7.shapes.add_textbox(Inches(2), Inches(2.0), Inches(6), Inches(4))
    tf = sponsor_box.text_frame
    tf.word_wrap = True
    for i, sp in enumerate(sponsors):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = sp
        run.font.size = Pt(24)
        run.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)
        p.space_after = Pt(18)

    # ========== Slide 8: Thank You ==========
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    fill8 = slide8.background.fill
    fill8.solid()
    fill8.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    add_centered_textbox(slide8, 'Thank You!', Inches(2.0), Pt(54),
                         RGBColor(0xFF, 0xFF, 0xFF), bold=True,
                         width=Inches(8), height=Inches(1.5))
    add_centered_textbox(slide8, 'See you next year!', Inches(4.0), Pt(24),
                         RGBColor(0xBD, 0xC3, 0xC7), bold=False,
                         width=Inches(8), height=Inches(0.8))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
