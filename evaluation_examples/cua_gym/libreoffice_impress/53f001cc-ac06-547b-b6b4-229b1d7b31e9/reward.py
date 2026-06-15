"""
Reward Script: Progress bar on every slide of a presentation
Task ID: impress_rp_042
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): All 10 slides have progress bar rectangles added
  Component 2 (0.35): Red accent bars (#E74C3C) with correct proportional widths
  Component 3 (0.20): Gray remainder bars (#E0E0E0) with correct complementary widths
  Component 4 (0.15): Bars positioned at slide bottom with correct height (0.15 inches)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_042'
FILE_PATH = os.path.join(WORKDIR, f'{TASK_ID}.pptx')

# Expected constants
ACCENT_COLOR = 'E74C3C'
GRAY_COLOR = 'E0E0E0'
BAR_HEIGHT_INCHES = 0.15
EMU_PER_INCH = 914400
EXPECTED_BAR_HEIGHT = int(BAR_HEIGHT_INCHES * EMU_PER_INCH)  # 137160 EMU
NUM_SLIDES = 10
# Tolerance for position/size comparisons (5% relative)
REL_TOL = 0.05


def get_fill_color(shape):
    """Get solid fill color as hex string, or None."""
    try:
        fill = shape.fill
        if fill.type == 1:  # solid fill
            return str(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def find_bar_shapes(slide):
    """Find rectangle shapes at the bottom of the slide that form the progress bar.
    Returns (red_shapes, gray_shapes) lists.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    red_shapes = []
    gray_shapes = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            color = get_fill_color(shape)
            if color == ACCENT_COLOR:
                red_shapes.append(shape)
            elif color == GRAY_COLOR:
                gray_shapes.append(shape)
    return red_shapes, gray_shapes


def approx_equal(a, b, rel_tol=REL_TOL):
    """Check if two values are approximately equal within relative tolerance."""
    if a == b:
        return True
    if a == 0 or b == 0:
        return abs(a - b) < rel_tol * max(abs(a), abs(b), 1)
    return abs(a - b) / max(abs(a), abs(b)) <= rel_tol


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slides = list(prs.slides)
    num_slides = len(slides)

    if num_slides != NUM_SLIDES:
        print(f"FAIL: Expected {NUM_SLIDES} slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: All 10 slides have at least one red progress bar rectangle (0.30 points)
    # Initial env has 0 such rectangles per slide; golden has 1+ per slide
    try:
        slides_with_bar = 0
        for i, slide in enumerate(slides):
            red_shapes, gray_shapes = find_bar_shapes(slide)
            if len(red_shapes) >= 1:
                slides_with_bar += 1
            else:
                print(f"  Slide {i+1}: no red progress bar found")

        if slides_with_bar == NUM_SLIDES:
            print(f"PASS: Component 1 — All {NUM_SLIDES} slides have red progress bar ({slides_with_bar}/{NUM_SLIDES}) (0.30 pts)")
            total_score += 0.30
        elif slides_with_bar > 0:
            partial = 0.30 * (slides_with_bar / NUM_SLIDES)
            print(f"PARTIAL: Component 1 — {slides_with_bar}/{NUM_SLIDES} slides have red progress bar ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — No slides have a red progress bar")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Red bars have correct proportional widths (0.35 points)
    # Slide N should have red bar width = slide_width * (N/10)
    try:
        correct_red_widths = 0
        for i, slide in enumerate(slides):
            slide_num = i + 1
            expected_red_width = int(slide_width * slide_num / NUM_SLIDES)
            red_shapes, _ = find_bar_shapes(slide)
            if red_shapes:
                # Find the red shape (should be one)
                red = red_shapes[0]
                actual_width = red.width
                if approx_equal(actual_width, expected_red_width):
                    correct_red_widths += 1
                else:
                    print(f"  Slide {slide_num}: red width {actual_width} != expected {expected_red_width}")
            else:
                print(f"  Slide {slide_num}: no red bar to check width")

        if correct_red_widths == NUM_SLIDES:
            print(f"PASS: Component 2 — All {NUM_SLIDES} red bars have correct proportional widths (0.35 pts)")
            total_score += 0.35
        elif correct_red_widths > 0:
            partial = 0.35 * (correct_red_widths / NUM_SLIDES)
            print(f"PARTIAL: Component 2 — {correct_red_widths}/{NUM_SLIDES} red bars correct ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No red bars have correct proportional widths")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Gray bars have correct color and complementary widths (0.20 points)
    # Slides 1-9 should have gray bar; slide 10 has 100% red so gray is optional/absent
    try:
        correct_gray = 0
        slides_needing_gray = NUM_SLIDES - 1  # slides 1-9
        for i, slide in enumerate(slides):
            slide_num = i + 1
            if slide_num == NUM_SLIDES:
                # Slide 10: 100% red, gray is optional (width 0 or absent)
                correct_gray += 0  # don't count, don't penalize
                continue
            expected_red_width = int(slide_width * slide_num / NUM_SLIDES)
            expected_gray_width = slide_width - expected_red_width
            _, gray_shapes = find_bar_shapes(slide)
            if gray_shapes:
                gray = gray_shapes[0]
                if approx_equal(gray.width, expected_gray_width):
                    correct_gray += 1
                else:
                    print(f"  Slide {slide_num}: gray width {gray.width} != expected {expected_gray_width}")
            else:
                print(f"  Slide {slide_num}: no gray bar found")

        if slides_needing_gray == 0:
            slides_needing_gray = 1  # avoid division by zero
        if correct_gray == NUM_SLIDES - 1:
            print(f"PASS: Component 3 — All {NUM_SLIDES - 1} gray bars have correct complementary widths (0.20 pts)")
            total_score += 0.20
        elif correct_gray > 0:
            partial = 0.20 * (correct_gray / slides_needing_gray)
            print(f"PARTIAL: Component 3 — {correct_gray}/{slides_needing_gray} gray bars correct ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No gray bars have correct complementary widths")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Bars positioned at slide bottom with correct height (0.15 points)
    # Expected: top position near bottom, height ~137160 EMU (0.15 inches)
    # Bar top should be approximately slide_height - bar_height
    try:
        correct_positions = 0
        expected_top = slide_height - EXPECTED_BAR_HEIGHT
        for i, slide in enumerate(slides):
            slide_num = i + 1
            red_shapes, gray_shapes = find_bar_shapes(slide)
            all_bar_shapes = red_shapes + gray_shapes
            if not all_bar_shapes:
                print(f"  Slide {slide_num}: no bar shapes to check position")
                continue

            position_errors = []
            for shape in all_bar_shapes:
                if not approx_equal(shape.top, expected_top):
                    position_errors.append(f"  Slide {slide_num}: bar top {shape.top} != expected {expected_top}")
                elif not approx_equal(shape.height, EXPECTED_BAR_HEIGHT):
                    position_errors.append(f"  Slide {slide_num}: bar height {shape.height} != expected {EXPECTED_BAR_HEIGHT}")
            if len(position_errors) == 0:
                correct_positions += 1
            else:
                for err in position_errors:
                    print(err)

        if correct_positions == NUM_SLIDES:
            print(f"PASS: Component 4 — All {NUM_SLIDES} slides have bars at correct position/height (0.15 pts)")
            total_score += 0.15
        elif correct_positions > 0:
            partial = 0.15 * (correct_positions / NUM_SLIDES)
            print(f"PARTIAL: Component 4 — {correct_positions}/{NUM_SLIDES} slides correct position ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — No slides have bars at correct position/height")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
persist_app_state()

if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(FILE_PATH)
