"""
Initial Setup: Create a 10-slide Course Module presentation with no progress bars.
Task ID: impress_rp_042
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_042'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard widescreen: 13.333 x 7.5 inches
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Course module content for 10 slides
    slide_content = [
        {
            "title": "Introduction to Data Analytics",
            "bullets": [
                "Course overview and learning objectives",
                "Prerequisites: basic statistics and spreadsheet skills",
                "Assessment breakdown: 40% assignments, 30% project, 30% exam",
                "Office hours: Tuesdays 2-4 PM, Room 312",
            ],
        },
        {
            "title": "Module 1: Data Collection Methods",
            "bullets": [
                "Primary vs secondary data sources",
                "Survey design principles and best practices",
                "Web scraping fundamentals with Python",
                "API integration for real-time data feeds",
                "Ethical considerations in data gathering",
            ],
        },
        {
            "title": "Module 2: Data Cleaning & Preparation",
            "bullets": [
                "Handling missing values: imputation strategies",
                "Outlier detection using IQR and Z-score methods",
                "Data type conversion and standardization",
                "Deduplication techniques for large datasets",
            ],
        },
        {
            "title": "Module 3: Exploratory Data Analysis",
            "bullets": [
                "Descriptive statistics: mean, median, mode, variance",
                "Correlation analysis and scatter plot interpretation",
                "Distribution visualization with histograms and box plots",
                "Identifying patterns and anomalies in datasets",
                "Case study: Retail sales trend analysis Q1-Q4 2024",
            ],
        },
        {
            "title": "Module 4: Statistical Inference",
            "bullets": [
                "Hypothesis testing: null vs alternative hypotheses",
                "Confidence intervals and margin of error",
                "Chi-square tests for categorical variables",
                "ANOVA for comparing multiple group means",
            ],
        },
        {
            "title": "Module 5: Regression Analysis",
            "bullets": [
                "Simple linear regression: assumptions and diagnostics",
                "Multiple regression with feature selection",
                "Logistic regression for binary classification",
                "R-squared interpretation and model comparison",
                "Practical exercise: predicting housing prices",
            ],
        },
        {
            "title": "Module 6: Data Visualization",
            "bullets": [
                "Choosing the right chart type for your data",
                "Color theory and accessibility in visualizations",
                "Interactive dashboards with Tableau and Power BI",
                "Storytelling with data: narrative structure",
            ],
        },
        {
            "title": "Module 7: Machine Learning Basics",
            "bullets": [
                "Supervised vs unsupervised learning paradigms",
                "Decision trees and random forest algorithms",
                "K-means clustering for customer segmentation",
                "Model evaluation: precision, recall, F1 score",
                "Cross-validation and preventing overfitting",
            ],
        },
        {
            "title": "Module 8: Working with Big Data",
            "bullets": [
                "Introduction to distributed computing frameworks",
                "SQL optimization for large-scale queries",
                "Cloud platforms: AWS, GCP, and Azure data services",
                "Data pipeline architecture and ETL processes",
            ],
        },
        {
            "title": "Final Project & Course Summary",
            "bullets": [
                "Project requirements: end-to-end analysis pipeline",
                "Deliverables: report, presentation, and code repository",
                "Submission deadline: March 28, 2025",
                "Key takeaways and recommended further reading",
                "Career paths in data analytics and data science",
            ],
        },
    ]

    for i, content in enumerate(slide_content):
        if i == 0:
            # Title slide layout
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = content["title"]
            slide.placeholders[1].text = "Dr. Elena Rodriguez | Spring 2025 Semester"
        else:
            # Title + Content layout
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = content["title"]
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for j, bullet in enumerate(content["bullets"]):
                if j == 0:
                    tf.paragraphs[0].text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
