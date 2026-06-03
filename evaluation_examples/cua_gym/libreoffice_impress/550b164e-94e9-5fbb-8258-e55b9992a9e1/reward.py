"""
Reward Script: Add merged title row to table on slide 5
Task ID: impress_tct_020
Domain: libreoffice_impress
Scoring:
  Component 1: Table has 6 rows (0.25)
  Component 2: Row 0 merged across all 4 columns (0.30)
  Component 3: Row 0 text is 'Sales Breakdown by Region' (0.25)
  Component 4: Row 1 contains original headers shifted down (0.20)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_020'


def persist_app_state(domain: str):
    """Save any unsaved GUI edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"PRECONDITION FAIL: Need at least 5 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Find the table on slide 5 (index 4)
    slide = prs.slides[4]
    table = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            break

    if table is None:
        print("PRECONDITION FAIL: No table found on slide 5")
        print("REWARD: 0.0")
        return 0.0

    num_rows = len(table.rows)
    num_cols = len(table.columns)
    print(f"INFO: Table dimensions: {num_rows} rows x {num_cols} cols")

    # Component 1: Table has 6 rows (0.25 points)
    # Initial has 5 rows; golden should have 6 (title row added)
    try:
        if num_rows == 6:
            print(f"PASS: Component 1 — Table has 6 rows (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — Expected 6 rows, found {num_rows}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Row 0 is merged across all 4 columns (0.30 points)
    # Check gridSpan attribute on the first cell of row 0
    try:
        if num_rows >= 6 and num_cols >= 4:
            tc0 = table.cell(0, 0)._tc
            grid_span = tc0.get('gridSpan')
            # Also verify the other cells have hMerge=1 (confirming horizontal merge)
            h_merge_count = 0
            for c in range(1, num_cols):
                tc = table.cell(0, c)._tc
                if tc.get('hMerge') == '1':
                    h_merge_count += 1

            if grid_span == '4' and h_merge_count == 3:
                print(f"PASS: Component 2 — Row 0 merged across 4 columns (gridSpan=4, 3 hMerge cells) (0.30 pts)")
                total_score += 0.30
            elif grid_span is not None and int(grid_span) >= 2:
                # Partial credit for partial merge
                print(f"PARTIAL: Component 2 — Row 0 gridSpan={grid_span}, hMerge cells={h_merge_count} (partial: 0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 2 — Row 0 not merged. gridSpan={grid_span}, hMerge cells={h_merge_count}")
        else:
            print(f"FAIL: Component 2 — Table too small for merge check ({num_rows}x{num_cols})")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Row 0 text is 'Sales Breakdown by Region' (0.25 points)
    try:
        if num_rows >= 6:
            row0_text = table.cell(0, 0).text.strip()
            if row0_text == 'Sales Breakdown by Region':
                print(f"PASS: Component 3 — Row 0 text matches exactly (0.25 pts)")
                total_score += 0.25
            elif 'sales breakdown' in row0_text.lower() and 'region' in row0_text.lower():
                # Partial credit for close match
                print(f"PARTIAL: Component 3 — Row 0 text close match: '{row0_text}' (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 3 — Expected 'Sales Breakdown by Region', found '{row0_text}'")
        else:
            print(f"FAIL: Component 3 — Not enough rows to check title")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Row 1 contains original headers (0.20 points)
    # The original headers were ['Region', 'Q1', 'Q2', 'Q3']
    # They should now be in row 1 (shifted down from row 0)
    try:
        if num_rows >= 6 and num_cols >= 4:
            expected_headers = ['Region', 'Q1', 'Q2', 'Q3']
            actual_headers = [table.cell(1, c).text.strip() for c in range(4)]
            if actual_headers == expected_headers:
                print(f"PASS: Component 4 — Row 1 headers correct: {actual_headers} (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 4 — Expected headers {expected_headers}, found {actual_headers}")
        else:
            print(f"FAIL: Component 4 — Not enough rows/cols for header check")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
