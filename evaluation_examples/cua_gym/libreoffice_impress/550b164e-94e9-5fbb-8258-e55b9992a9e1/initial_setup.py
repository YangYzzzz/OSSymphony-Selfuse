"""
Initial Setup: Create a 6-slide presentation with a table on slide 5.
Task ID: impress_tct_020
Domain: libreoffice_impress

Slide 5 has a 4-column x 5-row table with region sales data.
Row 1 = headers (Region, Q1, Q2, Q3), Rows 2-5 = data.
NO merged title row — that is the task for the agent to add.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_020'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_cell_text(cell, text, font_name="Calibri", font_size=Pt(14),
                  bold=False, color=None, alignment=None):
    """Helper to set cell text with formatting."""
    cell.text = ""
    tf = cell.text_frame
    p = tf.paragraphs[0]
    if alignment is not None:
        p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Regional Sales Report"
    slide1.placeholders[1].text = "FY 2025 Quarterly Performance Review"

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "This report covers quarterly sales across all four regions."
    p2 = body2.add_paragraph()
    p2.text = "Total revenue exceeded $2.4M in the fiscal year."
    p2.level = 0
    p3 = body2.add_paragraph()
    p3.text = "The West region showed the strongest growth trajectory."
    p3.level = 0
    p4 = body2.add_paragraph()
    p4.text = "East region maintained steady performance through all quarters."
    p4.level = 0

    # --- Slide 3: Market Trends ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Market Trends"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Consumer demand shifted towards digital channels in Q2."
    bp1 = body3.add_paragraph()
    bp1.text = "B2B segment grew 18% year-over-year in the North region."
    bp1.level = 0
    bp2 = body3.add_paragraph()
    bp2.text = "Retail partnerships expanded to 45 new locations."
    bp2.level = 0
    bp3 = body3.add_paragraph()
    bp3.text = "Product line diversification contributed to South region gains."
    bp3.level = 0

    # --- Slide 4: Key Metrics ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Key Performance Metrics"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Average deal size: $12,450"
    m1 = body4.add_paragraph()
    m1.text = "Customer retention rate: 92.3%"
    m1.level = 0
    m2 = body4.add_paragraph()
    m2.text = "New customer acquisition: 287 accounts"
    m2.level = 0
    m3 = body4.add_paragraph()
    m3.text = "Sales cycle length: 34 days (down from 41)"
    m3.level = 0

    # --- Slide 5: Sales Data Table (THE KEY SLIDE) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Add a title text box at top of slide
    title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Quarterly Sales Data"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Table: 5 rows x 4 columns (headers + 4 data rows)
    rows, cols = 5, 4
    table_shape = slide5.shapes.add_table(
        rows, cols, Inches(1.0), Inches(1.5), Inches(7.5), Inches(3.0)
    )
    table = table_shape.table

    # Set column widths
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(1.8)
    table.columns[3].width = Inches(1.9)

    # Row 1: Column headers
    headers = ['Region', 'Q1', 'Q2', 'Q3']
    header_color = RGBColor(0xFF, 0xFF, 0xFF)
    for c, h in enumerate(headers):
        set_cell_text(table.cell(0, c), h, bold=True, color=header_color,
                      alignment=PP_ALIGN.CENTER)
        # Header row background
        cell_fill = table.cell(0, c).fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Rows 2-5: Region data
    data = [
        ['North',  '$185,200', '$203,800', '$221,500'],
        ['South',  '$142,600', '$158,900', '$167,300'],
        ['East',   '$198,400', '$195,700', '$210,800'],
        ['West',   '$163,900', '$189,200', '$234,600'],
    ]
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            align = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            set_cell_text(table.cell(r, c), val, alignment=align)
            # Alternate row shading
            if r % 2 == 0:
                cell_fill = table.cell(r, c).fill
                cell_fill.solid()
                cell_fill.fore_color.rgb = RGBColor(0xE8, 0xEF, 0xF7)

    # --- Slide 6: Next Steps ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Next Steps"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Expand West region operations to capitalize on growth momentum."
    ns1 = body6.add_paragraph()
    ns1.text = "Launch targeted marketing campaign for South region in Q4."
    ns1.level = 0
    ns2 = body6.add_paragraph()
    ns2.text = "Hire 12 additional sales representatives across all regions."
    ns2.level = 0
    ns3 = body6.add_paragraph()
    ns3.text = "Schedule quarterly business reviews with top 20 accounts."
    ns3.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
