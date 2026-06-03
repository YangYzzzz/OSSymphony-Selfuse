"""
FINAL REWARD SCRIPT - SUCCESS
Task: Use Roman numerals (I, II, III) for pages 1–2, positioned top-center.
Generated: 2025-10-17 10:52:12
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation

# Reward script for verifying that Roman numerals (I, II, III …) are
# correctly placed on the first two slides, centred at the top.
# The script awards progressive points but returns exactly 1.0 only when
# ALL checks pass.

FILE_PATH = "/home/user/use_roman_numerals_i_ii_iii_for_pages_12_positioned_top_center.pptx"
EMU_PER_INCH = 914400  # PowerPoint internal units


def verify_roman_page_numbers(file_path: str) -> float:
    """Verify Roman numerals I and II are top-centered on slides 1 and 2.

    Scoring (progressive):
        • 0.25  per slide for correct numeral text
        • 0.25  per slide for correct top-centre positioning
        Maximum score = 1.0
    """

    total_score = 0.0
    max_score = 1.0

    # Check file exists (no points for existence – prerequisite)
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Unable to load presentation: {e}")
        return 0.0

    slide_width = prs.slide_width  # needed for centring calculation
    expected_numerals = {1: "I", 2: "II"}  # only pages 1–2 per instruction

    print(f"Loaded presentation with {len(prs.slides)} slides (width = {slide_width} EMU)")

    # Iterate through the first two slides (or fewer if presentation shorter)
    for idx, slide in enumerate(prs.slides, start=1):
        if idx > 2:
            break  # task only concerns pages 1–2

        expected_text = expected_numerals[idx]
        print(f"\nVerifying slide {idx} for numeral '{expected_text}' …")

        numeral_found = False
        correct_position = False

        # Search all text-bearing shapes
        for shape in slide.shapes:
            # Not all shapes have text
            if not hasattr(shape, "text"):
                continue

            text = shape.text.strip().upper()
            if text == expected_text:
                numeral_found = True

                # --- POSITION VALIDATION ---
                left = shape.left
                top = shape.top
                width = shape.width
                centre_x = left + width / 2

                # Allow up to 5 % horizontal tolerance from true centre
                horiz_tolerance = slide_width * 0.05
                # Must be within 1 inch of the top edge
                vertical_threshold = EMU_PER_INCH

                if abs(centre_x - slide_width / 2) <= horiz_tolerance and top <= vertical_threshold:
                    correct_position = True
                    print(f"  ✓ Numeral positioned correctly (top={top}, centre_x={centre_x})")
                else:
                    if abs(centre_x - slide_width / 2) > horiz_tolerance:
                        print(f"  ✗ Numeral not horizontally centred (centre_x diff={abs(centre_x - slide_width/2)})")
                    if top > vertical_threshold:
                        print(f"  ✗ Numeral too far from top (top={top})")

                break  # Stop once correct numeral located

        # ---- SCORING ----
        if numeral_found:
            print("  ✓ Correct numeral found")
            total_score += 0.25  # text correctness for this slide
        else:
            print("  ✗ Expected numeral not found on this slide")

        if numeral_found and correct_position:
            total_score += 0.25  # positioning correctness for this slide

    final_score = min(total_score, max_score)
    print(f"\nTotal score: {final_score} / {max_score}")
    return final_score


if __name__ == "__main__":
    reward = verify_roman_page_numbers(FILE_PATH)
    print(f"REWARD: {reward}")

