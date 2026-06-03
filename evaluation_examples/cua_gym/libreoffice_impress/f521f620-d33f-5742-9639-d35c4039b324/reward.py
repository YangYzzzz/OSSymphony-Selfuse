"""
FINAL REWARD SCRIPT - SUCCESS
Task: Insert a 7×5 empty table at the caret position.
Generated: 2025-10-17 11:08:16
Status: success
Model: azure-o3
Total Steps: 1
"""

from pptx import Presentation
import os

def verify_table_insert(file_path: str) -> float:
    """Verify that a 7×5 empty table was inserted at the caret position.

    Scoring (progressive):
      • 0.4 – Any table shape exists in the presentation
      • 0.4 – A table with exactly 7 columns × 5 rows is found
      • 0.2 – All cells in that 7×5 table are completely empty
      Maximum score = 1.0
    """
    max_score = 1.0
    score = 0.0

    print(f"Verifying file: {file_path}")

    # ---------- 1. File existence & load ----------
    if not os.path.exists(file_path):
        print("✗ File does not exist")
        return 0.0  # no points – task cannot be verified
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Failed to load PPTX: {e}")
        return 0.0  # loading failure yields 0

    # ---------- 2. Locate tables ----------
    table_found = False
    correct_dim_found = False
    empty_cells = False

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if not shape.has_table:
                continue  # skip non-table shapes

            table_found = True  # at least one table exists
            tbl = shape.table
            rows = len(tbl.rows)
            cols = len(tbl.columns)
            print(f"✓ Found table on slide {slide_idx}, shape {shape_idx}: {rows} rows × {cols} cols")

            # Check dimensions 5 rows × 7 columns (task specifies 7×5, meaning 7 columns and 5 rows)
            if rows == 5 and cols == 7:
                correct_dim_found = True

                # ---------- 3. Verify emptiness ----------
                all_empty = True
                for r in range(rows):
                    for c in range(cols):
                        text = tbl.cell(r, c).text_frame.text.strip()
                        if text:
                            all_empty = False
                            print(f"  ✗ Cell ({r},{c}) not empty: '{text}'")
                            break
                    if not all_empty:
                        break
                if all_empty:
                    empty_cells = True
                    print("  ✓ All cells are empty")

    # ---------- 4. Scoring ----------
    if table_found:
        score += 0.4
        print("✓ Table shape found (0.4)")
    else:
        print("✗ No table found (0)")

    if correct_dim_found:
        score += 0.4
        print("✓ Correct 7×5 table found (0.4)")
    else:
        print("✗ No table with 7×5 dimensions found (0)")

    if correct_dim_found and empty_cells:
        score += 0.2
        print("✓ Table cells are empty (0.2)")
    elif correct_dim_found:
        print("✗ Not all cells empty (0)")

    final_score = min(score, max_score)
    print(f"Final score: {final_score}")
    return final_score


if __name__ == "__main__":
    # Path to the presentation created by the agent
    path = "/home/user/insert_a_75_empty_table_at_the_caret_position.pptx"
    reward = verify_table_insert(path)
    print(f"REWARD: {reward}")
