"""
Reward Script: Create a summary slide listing all slide titles as bullet points
Task ID: impress_fix_066
Domain: libreoffice_impress
Scoring:
  Component 1 (0.2): Presentation has 13 slides (was 12)
  Component 2 (0.3): Slide 13 title is 'Summary'
  Component 3 (0.5): Slide 13 body lists all 12 original titles in order
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_066'

# The 12 original slide titles in order (from initial presentation)
EXPECTED_TITLES = [
    "Year in Review: 2025 Highlights",
    "Revenue Performance",
    "Customer Growth Metrics",
    "Product Development Milestones",
    "Marketing Campaign Results",
    "Employee Engagement & Culture",
    "Operational Efficiency Gains",
    "Strategic Partnerships",
    "Technology Infrastructure Updates",
    "Sustainability Initiatives",
    "Challenges & Lessons Learned",
    "2026 Goals & Roadmap",
]


def persist_app_state():
    """Save any unsaved LibreOffice edits before verification."""
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def get_all_text_shapes(slide):
    """Recursively get all text shapes including those inside groups."""
    def extract(shape):
        results = []
        if hasattr(shape, "text") and hasattr(shape, "text_frame"):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.extend(extract(sub))
        return results
    out = []
    for shape in slide.shapes:
        out.extend(extract(shape))
    return out


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Presentation has 13 slides (0.2 points)
    # Initial has 12, golden should have 13
    try:
        if num_slides == 13:
            print(f"PASS: Component 1 -- Slide count is 13 (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 1 -- Expected 13 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # If fewer than 13 slides, no summary slide to check
    if num_slides < 13:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Slide 13 (index 12) title is 'Summary' (0.3 points)
    try:
        last_slide = prs.slides[12]
        slide_title = ""
        if last_slide.shapes.title:
            slide_title = last_slide.shapes.title.text.strip()

        if slide_title.lower() == "summary":
            print(f"PASS: Component 2 -- Slide 13 title is '{slide_title}' (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 -- Expected title 'Summary', found '{slide_title}'")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Slide 13 body contains all 12 original titles in order (0.5 points)
    # Award partial credit: ~0.042 per correctly listed title in order
    try:
        last_slide = prs.slides[12]
        # Collect all non-title text from the summary slide
        body_texts = []
        for shape in get_all_text_shapes(last_slide):
            # Skip the title shape
            if last_slide.shapes.title and shape == last_slide.shapes.title:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    body_texts.append(text)

        print(f"  Summary slide body texts found: {body_texts}")

        # Check which expected titles appear in the body texts
        matched_count = 0
        body_lower = [bt.lower() for bt in body_texts]
        for expected_title in EXPECTED_TITLES:
            if any(expected_title.lower() in bt for bt in body_lower):
                matched_count += 1
            else:
                print(f"  MISS: Title '{expected_title}' not found in summary body")

        # Check ordering: titles should appear in the same order as in the presentation
        last_idx = -1
        order_violations = 0
        for expected_title in EXPECTED_TITLES:
            for j, bt in enumerate(body_lower):
                if expected_title.lower() in bt:
                    if j <= last_idx:
                        order_violations += 1
                    last_idx = j
                    break

        points_per_title = 0.5 / 12.0

        if matched_count == 12 and order_violations == 0:
            print(f"PASS: Component 3 -- All 12 titles found in order (0.5 pts)")
            total_score += 0.5
        elif matched_count == 12 and order_violations > 0:
            reduced = round(matched_count * points_per_title * 0.8, 3)
            print(f"PARTIAL: Component 3 -- All titles present but order incorrect ({reduced} pts)")
            if reduced > 0:
                total_score += reduced
        elif matched_count > 0:
            partial = round(matched_count * points_per_title, 3)
            print(f"PARTIAL: Component 3 -- {matched_count}/12 titles found ({partial} pts)")
            if partial > 0:
                total_score += partial
        else:
            print(f"FAIL: Component 3 -- No original titles found in summary slide body")

    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
