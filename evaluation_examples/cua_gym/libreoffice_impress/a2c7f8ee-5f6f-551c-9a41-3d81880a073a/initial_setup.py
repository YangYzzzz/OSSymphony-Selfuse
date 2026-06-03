"""
Initial Setup: Create a 12-slide Year End Review presentation
Task ID: impress_fix_066
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_066'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

SLIDE_TITLES = [
    "Year in Review: 2025 Highlights",
    "Revenue Performance",
    "Customer Growth Metrics",
    "Product Development Milestones",
    "Marketing Campaign Results",
    "Employee Engagement & Culture",
    "Operational Efficiency Gains",
    "Strategic Partnerships",
    "Technology Infrastructure Updates",
    "Sustainability Initiatives",
    "Challenges & Lessons Learned",
    "2026 Goals & Roadmap",
]

SLIDE_CONTENT = [
    "Welcome to the annual Year End Review.\nThis presentation covers key achievements, challenges, and strategic direction for our organization.",
    "Total revenue reached $48.7M, a 23% increase YoY.\nQ3 was the strongest quarter with $14.2M in bookings.\nEnterprise segment grew 31% driven by new contracts with Fortune 500 clients.",
    "Active users surpassed 1.2 million globally.\nCustomer retention rate improved to 94.3%.\nNet Promoter Score reached an all-time high of 72.",
    "Launched v3.0 platform with AI-powered analytics.\nShipped 47 feature releases across 4 product lines.\nReduced average bug resolution time from 5.2 to 2.1 days.",
    "Digital ad spend returned 4.8x ROAS.\nWebinar series attracted 15,000 registrations.\nBrand awareness in target markets increased by 18%.",
    "Employee satisfaction score: 4.2/5.0.\nHired 120 new team members across 8 departments.\nLaunched mentorship program pairing 85 junior and senior staff.",
    "Reduced deployment cycle from 2 weeks to 3 days.\nCloud infrastructure costs decreased 15% through optimization.\nAutomated 62% of manual QA processes.",
    "Signed 3 new technology partnership agreements.\nCo-developed integrations with Salesforce and ServiceNow.\nJoint go-to-market initiative generated $3.4M pipeline.",
    "Migrated 90% of workloads to multi-cloud architecture.\nAchieved 99.97% uptime across production services.\nCompleted SOC 2 Type II certification.",
    "Carbon footprint reduced by 22% year over year.\nTransitioned to 100% renewable energy in data centers.\nLaunched employee green commuting incentive program.",
    "Supply chain disruptions delayed Q2 hardware rollout.\nTalent acquisition in ML/AI roles remained competitive.\nKey takeaway: invest earlier in contingency planning.",
    "Target $62M revenue with focus on enterprise expansion.\nLaunch next-gen analytics platform in Q2.\nExpand into APAC and LATAM markets by end of year.",
]


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    for i, (title, content) in enumerate(zip(SLIDE_TITLES, SLIDE_CONTENT)):
        if i == 0:
            # Title slide for the first one
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title
            slide.placeholders[1].text = content
        else:
            # Title + Content layout for the rest
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title
            slide.placeholders[1].text = content

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Number of slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
