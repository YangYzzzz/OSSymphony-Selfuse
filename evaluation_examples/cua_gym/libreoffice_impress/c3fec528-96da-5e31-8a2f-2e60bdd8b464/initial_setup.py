"""
Initial Setup: Product launch presentation with 6 slides; slide 5 has two textboxes
Task ID: osworld_impress_textbox_fontsize_specific_005
Domain: libreoffice_impress
Note: Slide 5 textboxes use different font sizes than the target (54pt, 22pt) so the task is not pre-completed.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_textbox_fontsize_specific_005'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Use standard widescreen slide dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_layouts = prs.slide_layouts
    # Layout 0: Title Slide, 1: Title+Content, 5: Blank

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(slide_layouts[0])
    slide1.shapes.title.text = "NovaSpark Pro"
    slide1.placeholders[1].text = "The Next Generation Smart Home Platform\nProduct Launch 2025"

    # --- Slide 2: Market Opportunity ---
    slide2 = prs.slides.add_slide(slide_layouts[1])
    slide2.shapes.title.text = "Market Opportunity"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Smart home market reaching $180B by 2027"
    p2 = tf2.add_paragraph()
    p2.text = "42% of households plan to add smart devices this year"
    p3 = tf2.add_paragraph()
    p3.text = "Enterprise and residential segments both growing at 18% CAGR"
    p4 = tf2.add_paragraph()
    p4.text = "First-mover advantage in AI-integrated home automation"

    # --- Slide 3: Product Features ---
    slide3 = prs.slides.add_slide(slide_layouts[1])
    slide3.shapes.title.text = "Core Features"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "AI-Powered Automation — learns household routines in 7 days"
    p3b = tf3.add_paragraph()
    p3b.text = "Universal Compatibility — works with 3,000+ smart devices"
    p3c = tf3.add_paragraph()
    p3c.text = "Energy Savings Dashboard — average 31% reduction in utility bills"
    p3d = tf3.add_paragraph()
    p3d.text = "Enterprise Security — bank-grade 256-bit AES encryption"

    # --- Slide 4: Competitive Advantage ---
    slide4 = prs.slides.add_slide(slide_layouts[1])
    slide4.shapes.title.text = "Why NovaSpark Pro Wins"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "3x faster device onboarding than nearest competitor"
    p4b = tf4.add_paragraph()
    p4b.text = "Proprietary mesh protocol reduces latency by 85%"
    p4c = tf4.add_paragraph()
    p4c.text = "24/7 dedicated support with average 4-minute response time"
    p4d = tf4.add_paragraph()
    p4d.text = "5-year warranty — industry-leading commitment to quality"

    # --- Slide 5: Hero Statement (Key slide with two custom textboxes) ---
    slide5 = prs.slides.add_slide(slide_layouts[5])  # Blank layout

    # Background color for slide 5
    fill5 = slide5.background.fill
    fill5.solid()
    fill5.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # First textbox — hero statement (NOT 54pt — using 36pt so the task is not pre-completed)
    txBox1 = slide5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.0), Inches(2.5))
    tf_hero = txBox1.text_frame
    tf_hero.word_wrap = True
    p_hero = tf_hero.paragraphs[0]
    p_hero.alignment = PP_ALIGN.CENTER
    run_hero = p_hero.add_run()
    run_hero.text = "The Future of Smart Living Starts Here"
    run_hero.font.size = Pt(36)
    run_hero.font.bold = True
    run_hero.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Second textbox — supporting text (NOT 22pt — using 18pt so the task is not pre-completed)
    txBox2 = slide5.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10.0), Inches(1.8))
    tf_support = txBox2.text_frame
    tf_support.word_wrap = True
    p_support = tf_support.paragraphs[0]
    p_support.alignment = PP_ALIGN.CENTER
    run_support = p_support.add_run()
    run_support.text = "NovaSpark Pro transforms everyday spaces into intelligent environments that adapt to you."
    run_support.font.size = Pt(18)
    run_support.font.italic = True
    run_support.font.color.rgb = RGBColor(0xCC, 0xCC, 0xFF)

    # --- Slide 6: Call to Action ---
    slide6 = prs.slides.add_slide(slide_layouts[0])
    slide6.shapes.title.text = "Join the Revolution"
    slide6.placeholders[1].text = "Pre-order now at novaspark.io/pro\nEarly adopters receive 25% discount + free installation\nShipping begins Q2 2025"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
