"""
Reward Script: Set specific font sizes on slide 5 of a 6-slide presentation.
Task ID: osworld_impress_textbox_fontsize_specific_005
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.5): First textbox (TextBox 2) font size is 54pt
  - Component 2 (0.5): Second textbox (TextBox 3) font size is 22pt
  Total: 1.0
"""

import os

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_textbox_fontsize_specific_005'

# 1 pt = 12700 EMU in python-pptx font.size
PT_TO_EMU = 12700


def check_textbox_fontsize(shape, expected_pt):
    """
    Check that all non-empty runs in a textbox have the expected font size.
    Returns (runs_checked, runs_correct, failure_messages).
    """
    expected_emu = expected_pt * PT_TO_EMU
    runs_checked = 0
    runs_correct = 0
    failure_messages = []

    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if not (run.text or "").strip():
                continue  # skip empty/whitespace runs
            runs_checked += 1
            actual_size = run.font.size
            if actual_size == expected_emu:
                runs_correct += 1
            else:
                actual_pt = actual_size / PT_TO_EMU if actual_size else None
                preview = run.text[:30]
                failure_messages.append(
                    f"Run '{preview}' has size {actual_pt} pt, expected {expected_pt} pt"
                )

    return runs_checked, runs_correct, failure_messages


def verify_task(file_path):
    """
    Verify that slide 5 has:
      - First textbox (TextBox 2): all runs at 54pt
      - Second textbox (TextBox 3): all runs at 22pt
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"CRITICAL: Presentation has {len(prs.slides)} slides — expected at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide5 = prs.slides[4]  # 0-indexed, slide 5

    # Collect textbox shapes in order (excluding Title placeholders)
    textboxes = []
    for shape in slide5.shapes:
        if shape.has_text_frame and shape.name.startswith('TextBox'):
            textboxes.append(shape)

    if len(textboxes) < 2:
        print(f"FAIL: Expected at least 2 TextBox shapes on slide 5, found {len(textboxes)}")
        print("REWARD: 0.0")
        return 0.0

    first_tb = textboxes[0]   # TextBox 2 — hero statement
    second_tb = textboxes[1]  # TextBox 3 — supporting text

    # Component 1: First textbox (TextBox 2) font size == 54pt  (0.5 points)
    try:
        expected_pt = 54
        runs_checked, runs_correct, failures = check_textbox_fontsize(first_tb, expected_pt)

        if runs_checked > 0 and runs_checked == runs_correct:
            print(f"PASS Component 1: First textbox font size is {expected_pt}pt (shape: {first_tb.name!r})")
            total_score += 0.5
        elif runs_checked == 0:
            print(f"FAIL Component 1: No non-empty runs found in first textbox ({first_tb.name!r})")
        else:
            for msg in failures:
                print(f"FAIL Component 1: {msg}")

    except Exception as e:
        print(f"ERROR Component 1: Could not check first textbox font size: {e}")

    # Component 2: Second textbox (TextBox 3) font size == 22pt  (0.5 points)
    try:
        expected_pt = 22
        runs_checked, runs_correct, failures = check_textbox_fontsize(second_tb, expected_pt)

        if runs_checked > 0 and runs_checked == runs_correct:
            print(f"PASS Component 2: Second textbox font size is {expected_pt}pt (shape: {second_tb.name!r})")
            total_score += 0.5
        elif runs_checked == 0:
            print(f"FAIL Component 2: No non-empty runs found in second textbox ({second_tb.name!r})")
        else:
            for msg in failures:
                print(f"FAIL Component 2: {msg}")

    except Exception as e:
        print(f"ERROR Component 2: Could not check second textbox font size: {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
