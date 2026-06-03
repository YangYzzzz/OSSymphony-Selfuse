"""
Reward Script: Recipe presentation with 6 recipes, colored banners, index with hyperlinks
Task ID: impress_wf_029
Domain: libreoffice_impress
Scoring:
  C1 (0.15): File exists at Desktop/Cookbook.pptx with exactly 8 slides
  C2 (0.10): Slide 1 title contains "Family Favorites Cookbook"
  C3 (0.35): Slides 2-7 each have recipe structure (title, rounded rect, ingredients left, directions right)
  C4 (0.20): Slides 2-7 each have a colored banner at top with 6 distinct colors
  C5 (0.20): Slide 8 has index with 6 recipe names and hyperlinks to slides
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_029'

# Check file existence before importing heavy libs
_file_path = f'{WORKDIR}/Desktop/Cookbook.pptx'
if not os.path.exists(_file_path):
    print(f"File not found: {_file_path}")
    print("REWARD: 0.0")
    import sys
    sys.exit(0)

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn


def get_all_text_shapes(slide):
    """Get all shapes that have text frames, including in groups."""
    def extract(shape):
        results = []
        if hasattr(shape, "text") and hasattr(shape, "text_frame"):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.extend(extract(sub))
        return results
    out = []
    for shape in slide.shapes:
        out.extend(extract(shape))
    return out


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: File has exactly 8 slides (0.15 points)
    try:
        num_slides = len(prs.slides)
        if num_slides == 8:
            print(f"PASS: Component 1 — 8 slides found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — expected 8 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    if len(prs.slides) < 8:
        print(f"CRITICAL: Not enough slides ({len(prs.slides)}), cannot verify remaining components")
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Slide 1 title contains "Family Favorites Cookbook" (0.10 points)
    try:
        slide1 = prs.slides[0]
        all_text = ""
        for shape in get_all_text_shapes(slide1):
            all_text += " " + shape.text_frame.text
        all_text_lower = all_text.lower()
        if "family favorites cookbook" in all_text_lower:
            print(f"PASS: Component 2 — Title 'Family Favorites Cookbook' found on slide 1 (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 — 'Family Favorites Cookbook' not found on slide 1. Text: {all_text[:200]}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slides 2-7 recipe structure (0.35 points — ~0.058 per slide)
    # Each recipe slide needs: recipe title text, rounded rectangle with "Photo",
    # ingredients on left half, directions on right half
    try:
        recipe_structure_score = 0.0
        per_slide_pts = 0.35 / 6.0
        slide_width_mid = prs.slide_width / 2

        for idx in range(1, 7):  # slides index 1-6 = slides 2-7
            slide = prs.slides[idx]
            slide_num = idx + 1
            sub_score = 0.0
            sub_max = 4  # 4 checks per slide

            # Check a: Has a title-like text (recipe name in a text shape)
            text_shapes = get_all_text_shapes(slide)
            has_title = False
            for shape in text_shapes:
                text = shape.text_frame.text.strip()
                if text and len(text) > 3 and len(text) < 100 and "ingredient" not in text.lower() and "direction" not in text.lower() and "photo" not in text.lower():
                    has_title = True
                    break
            if has_title:
                sub_score += 1

            # Check b: Has a rounded rectangle with "Photo" text
            has_rounded_rect = False
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    try:
                        if shape.auto_shape_type is not None and shape.auto_shape_type == 5:  # ROUNDED_RECTANGLE
                            if hasattr(shape, 'text') and 'photo' in shape.text.lower():
                                has_rounded_rect = True
                                break
                    except Exception:
                        pass
                # Also check by name as fallback
                if 'rounded' in shape.name.lower() and 'rectangle' in shape.name.lower():
                    if hasattr(shape, 'text') and 'photo' in shape.text.lower():
                        has_rounded_rect = True
                        break
            if has_rounded_rect:
                sub_score += 1

            # Check c: Has ingredients text on left half
            has_ingredients_left = False
            for shape in text_shapes:
                text_lower = shape.text_frame.text.lower()
                if 'ingredient' in text_lower:
                    # Check position: left side means shape.left < slide_width_mid
                    center_x = shape.left + shape.width / 2
                    if center_x < slide_width_mid:
                        has_ingredients_left = True
                        break
            if has_ingredients_left:
                sub_score += 1

            # Check d: Has directions text on right half
            has_directions_right = False
            for shape in text_shapes:
                text_lower = shape.text_frame.text.lower()
                if 'direction' in text_lower or any(p.text.strip().startswith(('1.', '1)')) for p in shape.text_frame.paragraphs):
                    # Check position: right side means shape.left >= slide_width_mid
                    if shape.left >= slide_width_mid:
                        has_directions_right = True
                        break
            if has_directions_right:
                sub_score += 1

            slide_frac = (sub_score / sub_max) * per_slide_pts
            recipe_structure_score += slide_frac

            if sub_score == sub_max:
                print(f"  Slide {slide_num}: PASS all 4 checks")
            else:
                print(f"  Slide {slide_num}: {sub_score}/{sub_max} checks passed")

        recipe_structure_score = round(recipe_structure_score, 4)
        if recipe_structure_score > 0:
            print(f"PASS: Component 3 — Recipe structure score: {recipe_structure_score:.4f}/0.35 ({recipe_structure_score:.4f} pts)")
            total_score += recipe_structure_score
        else:
            print(f"FAIL: Component 3 — No recipe structure found on slides 2-7")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Colored banners at top of slides 2-7, 6 distinct colors (0.20 points)
    try:
        banner_colors = []
        banners_found = 0

        for idx in range(1, 7):  # slides 2-7
            slide = prs.slides[idx]
            slide_num = idx + 1
            found_banner = False

            for shape in slide.shapes:
                # Banner: a rectangle at the top of the slide (top position near 0)
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    # Check if it spans wide and is near top
                    if shape.top < 1000000 and shape.width > prs.slide_width * 0.7:
                        try:
                            fill = shape.fill
                            if fill.type is not None and fill.type == 1:  # SOLID
                                color_rgb = str(fill.fore_color.rgb)
                                banner_colors.append(color_rgb)
                                banners_found += 1
                                found_banner = True
                                print(f"  Slide {slide_num}: banner color {color_rgb}")
                                break
                        except Exception:
                            pass

            if not found_banner:
                print(f"  Slide {slide_num}: no banner found")

        # Score: proportional to banners found + bonus for distinct colors
        banner_pts = 0.0
        if banners_found > 0:
            # 0.10 for having banners on all 6 slides
            banner_pts += 0.10 * (banners_found / 6.0)
            # 0.10 for having 6 distinct colors
            unique_colors = len(set(banner_colors))
            if unique_colors >= 6:
                banner_pts += 0.10
                print(f"  All {unique_colors} banner colors are distinct")
            elif unique_colors >= 3:
                banner_pts += 0.05
                print(f"  Only {unique_colors}/6 distinct banner colors")
            else:
                print(f"  Only {unique_colors}/6 distinct banner colors")

        banner_pts = round(banner_pts, 4)
        if banner_pts > 0:
            print(f"PASS: Component 4 — Banner score: {banner_pts:.4f}/0.20 ({banner_pts:.4f} pts)")
            total_score += banner_pts
        else:
            print(f"FAIL: Component 4 — No colored banners found on slides 2-7")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slide 8 index with 6 recipe names and hyperlinks (0.20 points)
    try:
        slide8 = prs.slides[7]
        text_shapes = get_all_text_shapes(slide8)

        # Collect all text on slide 8
        all_texts = []
        for shape in text_shapes:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    all_texts.append(t)

        # Check for recipe names — collect recipe titles from slides 2-7
        recipe_titles = []
        for idx in range(1, 7):
            slide = prs.slides[idx]
            for shape in get_all_text_shapes(slide):
                text = shape.text_frame.text.strip()
                # The title is typically the first meaningful non-banner, non-ingredients text
                if text and len(text) > 3 and len(text) < 100:
                    if "ingredient" not in text.lower() and "direction" not in text.lower() and "photo" not in text.lower():
                        recipe_titles.append(text.split('\n')[0].strip())
                        break

        # Check how many recipe titles appear in slide 8 text
        slide8_full_text = " ".join(all_texts).lower()
        recipes_found = 0
        for title in recipe_titles:
            # Check if the core recipe name appears (allow partial match)
            title_words = title.lower().strip()
            if title_words in slide8_full_text:
                recipes_found += 1

        # Check for hyperlinks on slide 8
        hyperlink_count = 0
        for shape in text_shapes:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    # Check XML for internal slide hyperlinks
                    rPr = run._r.find(qn('a:rPr'))
                    if rPr is not None:
                        hlinkClick = rPr.find(qn('a:hlinkClick'))
                        if hlinkClick is not None:
                            action = hlinkClick.get('action', '')
                            if 'hlinksldjump' in action or hlinkClick.get(qn('r:id'), ''):
                                hyperlink_count += 1

        index_pts = 0.0
        # 0.10 for recipe names present
        if recipes_found >= 6:
            index_pts += 0.10
            print(f"  All 6 recipe names found on slide 8")
        elif recipes_found > 0:
            index_pts += 0.10 * (recipes_found / 6.0)
            print(f"  {recipes_found}/6 recipe names found on slide 8")
        else:
            print(f"  No recipe names found on slide 8")

        # 0.10 for hyperlinks
        if hyperlink_count >= 6:
            index_pts += 0.10
            print(f"  {hyperlink_count} slide hyperlinks found on slide 8")
        elif hyperlink_count > 0:
            index_pts += 0.10 * (hyperlink_count / 6.0)
            print(f"  {hyperlink_count}/6 hyperlinks found on slide 8")
        else:
            print(f"  No hyperlinks found on slide 8")

        index_pts = round(index_pts, 4)
        if index_pts > 0:
            print(f"PASS: Component 5 — Index score: {index_pts:.4f}/0.20 ({index_pts:.4f} pts)")
            total_score += index_pts
        else:
            print(f"FAIL: Component 5 — Slide 8 has no recipe index")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point — file existence already checked at top of script
verify_task(_file_path)
