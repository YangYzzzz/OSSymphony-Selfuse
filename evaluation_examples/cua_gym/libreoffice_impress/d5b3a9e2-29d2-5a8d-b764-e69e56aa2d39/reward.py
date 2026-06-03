"""
Reward Script: Duplicate slide 4 and place as slide 5 with title 'Alternative Approach'
Task ID: impress_gf3_005
Domain: libreoffice_impress
Scoring:
  Component 1: Total slide count is 8 (0.2 pts)
  Component 2: Slide 5 title text is 'Alternative Approach' (0.3 pts)
  Component 3: Original slide 4 title unchanged as 'Recommended Approach' (0.2 pts)
  Component 4: Slide 5 duplicates slide 4 non-title content (0.2 pts)
  Component 5: Former slides 5-7 shifted to 6-8 correctly (0.1 pts)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_005'


def get_title_text(slide):
    """Get the title-like text from a slide (first text shape with short text)."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            # Title is typically the first short text (the heading)
            # Look for TextBox named 'TextBox 2' which is the title in this pptx
            if shape.name == 'TextBox 2' and text:
                return text
    return None


def get_non_title_texts(slide):
    """Get all non-title text content from a slide, sorted for comparison."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            if shape.name == 'TextBox 2':
                continue  # skip title
            text = shape.text_frame.text.strip()
            if text:
                texts.append(text)
    texts.sort()
    return texts


def get_non_title_shape_info(slide):
    """Get shape types and counts excluding the title textbox, for structural comparison."""
    info = []
    for shape in slide.shapes:
        if shape.name == 'TextBox 2':
            continue
        if hasattr(shape, 'name') and shape.name == 'Title 1':
            continue  # skip empty placeholder title
        info.append(str(shape.shape_type))
    info.sort()
    return info


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Slide count is 8 (0.2 points)
    # Initial has 7 slides; golden has 8 (one duplicate added)
    try:
        if num_slides == 8:
            print(f"PASS: Component 1 — Slide count is 8 (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 1 — Expected 8 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # If fewer than 5 slides, remaining checks don't make sense
    if num_slides < 5:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Slide 5 title is 'Alternative Approach' (0.3 points)
    # This is the core task change — initial slide 5 has 'Timeline & Milestones'
    try:
        slide5 = prs.slides[4]  # 0-indexed
        slide5_title = get_title_text(slide5)
        if slide5_title == 'Alternative Approach':
            print(f"PASS: Component 2 — Slide 5 title is 'Alternative Approach' (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 — Expected slide 5 title 'Alternative Approach', found: {repr(slide5_title)}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Original slide 4 title unchanged as 'Recommended Approach' (0.2 points)
    # Both initial and golden have this — BUT we combine with Component 1 (slide count == 8)
    # to ensure this only passes when the duplication happened.
    # On initial_env: slide count != 8, so Component 1 already fails.
    # This component checks slide 4 still says 'Recommended Approach' AND slide 5 says
    # 'Alternative Approach' (anchored to task change).
    try:
        slide4 = prs.slides[3]  # 0-indexed
        slide4_title = get_title_text(slide4)
        slide5_title_check = get_title_text(prs.slides[4]) if num_slides >= 5 else None
        if slide4_title == 'Recommended Approach' and slide5_title_check == 'Alternative Approach':
            print(f"PASS: Component 3 — Slide 4 title is 'Recommended Approach' AND slide 5 is 'Alternative Approach' (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 3 — Slide 4 title: {repr(slide4_title)}, Slide 5 title: {repr(slide5_title_check)}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 5 duplicates slide 4 non-title content (0.2 points)
    # The duplicate should have same shapes and non-title text as slide 4
    # On initial_env, slide 5 is 'Timeline & Milestones' — different content entirely
    try:
        slide4 = prs.slides[3]
        slide5 = prs.slides[4]

        s4_texts = get_non_title_texts(slide4)
        s5_texts = get_non_title_texts(slide5)

        s4_shapes = get_non_title_shape_info(slide4)
        s5_shapes = get_non_title_shape_info(slide5)

        texts_match = s4_texts == s5_texts
        shapes_match = s4_shapes == s5_shapes

        if texts_match and shapes_match:
            print(f"PASS: Component 4 — Slide 5 non-title content matches slide 4 (0.2 pts)")
            total_score += 0.2
        else:
            if not texts_match:
                print(f"FAIL: Component 4 — Non-title texts differ. Slide 4: {len(s4_texts)} items, Slide 5: {len(s5_texts)} items")
            if not shapes_match:
                print(f"FAIL: Component 4 — Shape types differ. Slide 4: {s4_shapes}, Slide 5: {s5_shapes}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Former slides 5-7 shifted to 6-8 (0.1 points)
    # On initial, slide 5='Timeline & Milestones', slide 6='Budget Breakdown', slide 7='Next Steps'
    # On golden, slide 6='Timeline & Milestones', slide 7='Budget Breakdown', slide 8='Next Steps'
    # On initial_env, slide 6 (if exists) would be 'Budget Breakdown', not 'Timeline & Milestones'
    try:
        if num_slides >= 8:
            s6_title = get_title_text(prs.slides[5])
            s7_title = get_title_text(prs.slides[6])
            s8_title = get_title_text(prs.slides[7])

            expected = ['Timeline & Milestones', 'Budget Breakdown', 'Next Steps']
            actual = [s6_title, s7_title, s8_title]

            if actual == expected:
                print(f"PASS: Component 5 — Slides 6-8 correctly shifted (0.1 pts)")
                total_score += 0.1
            else:
                print(f"FAIL: Component 5 — Expected slides 6-8 titles {expected}, found {actual}")
        else:
            print(f"FAIL: Component 5 — Need 8 slides for shift check, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice (best-effort save before verification)
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
