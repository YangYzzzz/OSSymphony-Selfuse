"""
Initial Setup: Create a 7-slide Project Proposal presentation
Task ID: impress_gf3_005
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_005'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_shape(shape, text, font_size=14, bold=False, color=None, alignment=None):
    """Helper to set text on a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ============ Slide 1: Title Slide ============
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Project Proposal"
    slide1.placeholders[1].text = "Digital Transformation Initiative — Q2 2025\nPrepared by Elena Rodriguez, VP of Strategy"

    # ============ Slide 2: Project Overview ============
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title text box
    title2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
    add_text_to_shape(title2, "Project Overview", font_size=28, bold=True,
                      color=RGBColor(0x1F, 0x49, 0x7D))
    # Body text
    body2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(5))
    tf2 = body2.text_frame
    tf2.word_wrap = True
    items = [
        "Our organization faces increasing competition from digital-native companies in the fintech sector.",
        "Customer retention rates declined 12% year-over-year, primarily due to outdated mobile experiences.",
        "Internal processes still rely on legacy systems deployed in 2014, resulting in 23% higher operational costs.",
        "This proposal outlines a comprehensive 18-month transformation plan targeting three core areas:",
        "  1. Customer-facing digital products (mobile app redesign, web portal upgrade)",
        "  2. Internal workflow automation (AI-powered document processing, automated reporting)",
        "  3. Data infrastructure modernization (cloud migration, real-time analytics pipeline)",
    ]
    for i, item in enumerate(items):
        if i == 0:
            tf2.paragraphs[0].text = item
        else:
            p = tf2.add_paragraph()
            p.text = item
        para = tf2.paragraphs[i]
        for run in para.runs:
            run.font.size = Pt(14)

    # ============ Slide 3: Market Analysis ============
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    title3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
    add_text_to_shape(title3, "Market Analysis", font_size=28, bold=True,
                      color=RGBColor(0x1F, 0x49, 0x7D))

    body3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5), Inches(5))
    tf3 = body3.text_frame
    tf3.word_wrap = True
    lines3 = [
        "Key Market Findings:",
        "• 78% of competitors launched mobile-first platforms in the past 2 years",
        "• Average customer expects sub-3-second page load times",
        "• AI-driven personalization increases conversion rates by 34%",
        "• Cloud-native companies report 40% lower infrastructure costs",
        "",
        "Target Segments:",
        "• Enterprise clients (revenue > $50M) — 42% of current portfolio",
        "• Mid-market SaaS companies — fastest growing at 28% CAGR",
        "• Government & regulatory bodies — high-value, long-term contracts",
    ]
    for i, line in enumerate(lines3):
        if i == 0:
            tf3.paragraphs[0].text = line
            tf3.paragraphs[0].runs[0].font.bold = True
        else:
            p = tf3.add_paragraph()
            p.text = line
        for run in tf3.paragraphs[i].runs:
            run.font.size = Pt(13)

    # Stats box on the right
    stats3 = slide3.shapes.add_textbox(Inches(7), Inches(1.8), Inches(5), Inches(4))
    tf_stats = stats3.text_frame
    tf_stats.word_wrap = True
    stat_items = [
        ("$4.2B", "Total addressable market"),
        ("18%", "Expected market growth (CAGR)"),
        ("67%", "Customer willingness to switch for better UX"),
        ("$850K", "Average annual contract value target"),
    ]
    for i, (val, label) in enumerate(stat_items):
        if i == 0:
            p = tf_stats.paragraphs[0]
        else:
            p = tf_stats.add_paragraph()
        p.text = f"{val} — {label}"
        for run in p.runs:
            run.font.size = Pt(14)

    # ============ Slide 4: Recommended Approach (KEY SLIDE) ============
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    title4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
    add_text_to_shape(title4, "Recommended Approach", font_size=28, bold=True,
                      color=RGBColor(0x1F, 0x49, 0x7D))

    # Three diagram shapes (rounded rectangles representing phases)
    phase_colors = [
        RGBColor(0x2E, 0x75, 0xB6),  # Phase 1 blue
        RGBColor(0x44, 0x72, 0xC4),  # Phase 2 medium blue
        RGBColor(0x5B, 0x9B, 0xD5),  # Phase 3 light blue
    ]
    phase_labels = ["Phase 1:\nDiscovery &\nPlanning", "Phase 2:\nDevelopment &\nTesting", "Phase 3:\nDeployment &\nScaling"]
    for i in range(3):
        shape = slide4.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.0 + i * 3.8), Inches(2.0),
            Inches(3.0), Inches(2.0)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = phase_colors[i]
        shape.line.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = phase_labels[i]
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Two text boxes below the diagram
    tb1 = slide4.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(5), Inches(2))
    tf_tb1 = tb1.text_frame
    tf_tb1.word_wrap = True
    tf_tb1.paragraphs[0].text = "Timeline: 18 months total"
    p_tb1 = tf_tb1.add_paragraph()
    p_tb1.text = "• Phase 1: Months 1-4 (stakeholder interviews, tech audit, roadmap)"
    p_tb1a = tf_tb1.add_paragraph()
    p_tb1a.text = "• Phase 2: Months 5-12 (agile sprints, MVP releases, user testing)"
    p_tb1b = tf_tb1.add_paragraph()
    p_tb1b.text = "• Phase 3: Months 13-18 (production rollout, training, optimization)"
    for para in tf_tb1.paragraphs:
        for run in para.runs:
            run.font.size = Pt(12)

    tb2 = slide4.shapes.add_textbox(Inches(7), Inches(4.5), Inches(5), Inches(2))
    tf_tb2 = tb2.text_frame
    tf_tb2.word_wrap = True
    tf_tb2.paragraphs[0].text = "Investment Required: $3.2M"
    tf_tb2.paragraphs[0].runs[0].font.bold = True
    p_tb2 = tf_tb2.add_paragraph()
    p_tb2.text = "Expected ROI: 240% over 3 years"
    p_tb2a = tf_tb2.add_paragraph()
    p_tb2a.text = "Break-even: Month 14"
    for para in tf_tb2.paragraphs:
        for run in para.runs:
            run.font.size = Pt(12)

    # ============ Slide 5: Timeline & Milestones ============
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    title5 = slide5.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
    add_text_to_shape(title5, "Timeline & Milestones", font_size=28, bold=True,
                      color=RGBColor(0x1F, 0x49, 0x7D))

    body5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(5))
    tf5 = body5.text_frame
    tf5.word_wrap = True
    milestones = [
        "Q2 2025 — Kickoff & Discovery Phase",
        "  • Stakeholder interviews completed (45 sessions across 8 departments)",
        "  • Technical architecture review and vendor selection",
        "Q3 2025 — Foundation Sprint",
        "  • Cloud infrastructure provisioned (AWS + Azure hybrid)",
        "  • Core API framework deployed to staging",
        "Q4 2025 — MVP & User Testing",
        "  • Mobile app beta released to 500 internal users",
        "  • Automated document processing pilot with Legal team",
        "Q1 2026 — Production Rollout",
        "  • Full mobile app launch to all 12,000 customers",
        "  • Legacy system decommissioning begins",
        "Q2 2026 — Optimization & Scaling",
        "  • Performance tuning based on production metrics",
        "  • Phase 2 feature roadmap finalized",
    ]
    for i, item in enumerate(milestones):
        if i == 0:
            tf5.paragraphs[0].text = item
        else:
            p = tf5.add_paragraph()
            p.text = item
        for run in tf5.paragraphs[i].runs:
            run.font.size = Pt(13)
            if not item.startswith("  "):
                run.font.bold = True

    # ============ Slide 6: Budget Breakdown ============
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    title6 = slide6.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
    add_text_to_shape(title6, "Budget Breakdown", font_size=28, bold=True,
                      color=RGBColor(0x1F, 0x49, 0x7D))

    # Budget table
    table_shape = slide6.shapes.add_table(6, 4, Inches(0.8), Inches(1.5), Inches(11), Inches(3.5))
    table = table_shape.table
    headers = ["Category", "Phase 1", "Phase 2", "Phase 3"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(13)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    budget_data = [
        ["Engineering & Development", "$280,000", "$920,000", "$380,000"],
        ["Cloud Infrastructure", "$45,000", "$185,000", "$210,000"],
        ["UX Design & Research", "$120,000", "$95,000", "$40,000"],
        ["Training & Change Management", "$15,000", "$60,000", "$145,000"],
        ["Contingency (15%)", "$69,000", "$189,000", "$116,250"],
    ]
    for r, row in enumerate(budget_data, 1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(12)

    # Total line below table
    total_box = slide6.shapes.add_textbox(Inches(0.8), Inches(5.3), Inches(11), Inches(0.6))
    add_text_to_shape(total_box, "Total Investment: $3,200,000 (including 15% contingency reserve)",
                      font_size=14, bold=True, color=RGBColor(0x1F, 0x49, 0x7D))

    # ============ Slide 7: Next Steps ============
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    title7 = slide7.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
    add_text_to_shape(title7, "Next Steps", font_size=28, bold=True,
                      color=RGBColor(0x1F, 0x49, 0x7D))

    body7 = slide7.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(5))
    tf7 = body7.text_frame
    tf7.word_wrap = True
    next_steps = [
        "Immediate Actions (Next 2 Weeks):",
        "  1. Executive steering committee approval — present to Board on April 28",
        "  2. Finalize vendor shortlist for cloud migration partner",
        "  3. Begin hiring for 4 senior engineering positions",
        "",
        "Key Decision Points:",
        "  • Build vs. buy for mobile app framework (recommendation: hybrid approach)",
        "  • Single-cloud vs. multi-cloud strategy (recommendation: AWS primary + Azure DR)",
        "  • Internal team vs. external consultants for UX research",
        "",
        "Contact: Elena Rodriguez (e.rodriguez@company.com) | ext. 4521",
    ]
    for i, item in enumerate(next_steps):
        if i == 0:
            tf7.paragraphs[0].text = item
        else:
            p = tf7.add_paragraph()
            p.text = item
        for run in tf7.paragraphs[i].runs:
            run.font.size = Pt(13)
            if item.endswith(":"):
                run.font.bold = True

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
