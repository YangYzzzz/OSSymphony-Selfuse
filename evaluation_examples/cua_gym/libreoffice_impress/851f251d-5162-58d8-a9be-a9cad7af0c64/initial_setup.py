"""
Initial Setup: Create a 7-slide Revenue Pitch presentation with empty content on slide 4.
Task ID: impress_sales_022
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_022'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_body(slide, title_text, body_lines):
    """Set title and add body text lines to a slide with title+content layout."""
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide (layout 0) ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Revenue Pitch 2024"
    slide1.placeholders[1].text = "Acme Global Solutions\nAnnual Performance Review"

    # --- Slide 2: Company Overview (layout 1 = Title + Content) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide2, "Company Overview", [
        "Founded in 2015 with headquarters in San Francisco",
        "Over 1,200 employees across 8 global offices",
        "Serving Fortune 500 clients in 14 countries",
        "Annual recurring revenue exceeded $12M in 2023",
        "Named Top 50 Fastest Growing SaaS Companies by Forbes",
    ])

    # --- Slide 3: Market Analysis (layout 1) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide3, "Market Analysis", [
        "Total addressable market projected at $48B by 2026",
        "Cloud adoption rate growing 22% year-over-year",
        "Key competitors: TechNova, DataSphere, CloudVista",
        "Our market share increased from 3.2% to 5.8% in 2024",
        "Enterprise segment shows strongest demand signals",
    ])

    # --- Slide 4: Our Growth Story (layout 1 — empty content area, NO chart) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Our Growth Story"
    # Leave content placeholder empty — the task is to add a bar chart here

    # --- Slide 5: Client Portfolio (layout 1) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide5, "Client Portfolio", [
        "Meridian Healthcare — EHR platform migration ($1.8M deal)",
        "Pinnacle Financial Group — Real-time analytics dashboard",
        "Atlas Logistics — Supply chain optimization suite",
        "NovaTech Industries — IoT sensor data pipeline",
        "92% client retention rate over the past 3 years",
    ])

    # --- Slide 6: Strategic Initiatives (layout 1) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide6, "Strategic Initiatives", [
        "Launch AI-powered predictive analytics module in Q2",
        "Expand APAC sales team by 40% through new hires",
        "Achieve SOC 2 Type II and ISO 27001 certifications",
        "Establish strategic partnership with AWS Marketplace",
        "Target $18M ARR by end of fiscal year 2025",
    ])

    # --- Slide 7: Next Steps (layout 1) ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide7, "Next Steps", [
        "Schedule follow-up meeting with investment committee",
        "Distribute detailed financial projections by March 15",
        "Arrange product demo for technical due diligence",
        "Finalize term sheet discussions by end of Q1",
        "Contact: Sarah Chen, VP of Strategy — sarah@acmeglobal.com",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
