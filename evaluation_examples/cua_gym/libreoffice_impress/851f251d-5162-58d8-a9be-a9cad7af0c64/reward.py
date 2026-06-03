"""
Reward Script: Create a bar chart on slide 4 with quarterly revenue data
Task ID: impress_sales_022
Domain: libreoffice_impress
Scoring:
  - Component 1: Chart exists on slide 4 (0.25)
  - Component 2: Chart type is bar/column (0.15)
  - Component 3: Chart title is 'Revenue Growth 2024' (0.20)
  - Component 4: Categories are Q1, Q2, Q3, Q4 (0.20)
  - Component 5: Data values match [2.1, 2.8, 3.5, 4.2] (0.20)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_022'


def persist_app_state():
    """Save any unsaved LibreOffice edits via Ctrl+S."""
    try:
        os.environ["DISPLAY"] = ":0"
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Presentation has at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide4 = prs.slides[3]  # 0-indexed, slide 4

    # Find chart shape(s) on slide 4
    chart_shapes = [s for s in slide4.shapes if s.has_chart]

    # Component 1: Chart exists on slide 4 (0.25 points)
    try:
        if len(chart_shapes) > 0:
            print(f"PASS: Component 1 -- Chart found on slide 4 ({len(chart_shapes)} chart(s)) (0.25 pts)")
            total_score += 0.25
        else:
            print("FAIL: Component 1 -- No chart found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    if len(chart_shapes) == 0:
        # No chart means nothing else to check
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    chart = chart_shapes[0].chart

    # Component 2: Chart type is bar or column (0.15 points)
    try:
        from pptx.enum.chart import XL_CHART_TYPE
        # Accept any bar or column variant
        bar_column_types = {
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            XL_CHART_TYPE.COLUMN_STACKED,
            XL_CHART_TYPE.COLUMN_STACKED_100,
            XL_CHART_TYPE.BAR_CLUSTERED,
            XL_CHART_TYPE.BAR_STACKED,
            XL_CHART_TYPE.BAR_STACKED_100,
        }
        ct = chart.chart_type
        if ct in bar_column_types:
            print(f"PASS: Component 2 -- Chart type is bar/column ({ct}) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 -- Chart type is {ct}, expected bar or column")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Chart title is 'Revenue Growth 2024' (0.20 points)
    try:
        if chart.has_title:
            title_text = chart.chart_title.text_frame.text.strip()
            if title_text == "Revenue Growth 2024":
                print(f"PASS: Component 3 -- Chart title is 'Revenue Growth 2024' (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 3 -- Chart title is '{title_text}', expected 'Revenue Growth 2024'")
        else:
            print("FAIL: Component 3 -- Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Categories are Q1, Q2, Q3, Q4 (0.20 points)
    try:
        plot = chart.plots[0]
        cats = list(plot.categories)
        expected_cats = ['Q1', 'Q2', 'Q3', 'Q4']
        if cats == expected_cats:
            print(f"PASS: Component 4 -- Categories match {expected_cats} (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 4 -- Categories are {cats}, expected {expected_cats}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Data values match [2.1, 2.8, 3.5, 4.2] (0.20 points)
    try:
        if len(chart.series) > 0:
            actual_values = list(chart.series[0].values)
            expected_values = [2.1, 2.8, 3.5, 4.2]
            # Use tolerance for float comparison
            values_match = len(actual_values) == len(expected_values) and all(
                abs(a - e) < 0.01 for a, e in zip(actual_values, expected_values)
            )
            if values_match:
                print(f"PASS: Component 5 -- Values match {expected_values} (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 5 -- Values are {actual_values}, expected {expected_values}")
        else:
            print("FAIL: Component 5 -- No data series found in chart")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
