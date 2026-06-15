"""
Reward Script: Add agenda slide with hyperlinked items to Startup_Pitch.pptx
Task ID: impress_ps_042
Domain: libreoffice_impress
Scoring:
  Component 1: Slide count is 10 (0.15)
  Component 2: Slide 2 has 'Agenda' title (0.15)
  Component 3: 7 agenda items with correct text (0.30)
  Component 4: Hyperlinks to correct slides (0.25)
  Component 5: Font formatting - 18pt, blue, underlined (0.15)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_042'

EXPECTED_ITEMS = [
    'Problem Statement',
    'Our Solution',
    'Market Opportunity',
    'Competitive Landscape',
    'Financial Projections',
    'Our Team',
    'The Ask',
]

# Each item at index i should link to slide (i+3) in the presentation
# (item 0 -> slide 3, item 1 -> slide 4, ..., item 6 -> slide 9)


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.oxml.ns import qn
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide count is 10 (0.15 points)
    # Initial has 9 slides; golden should have 10 (agenda inserted at position 2)
    try:
        num_slides = len(prs.slides)
        if num_slides == 10:
            print(f"PASS: Component 1 -- Slide count is 10 (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 -- Expected 10 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Slide 2 has 'Agenda' title (0.15 points)
    # Initial slide 2 is 'Problem Statement'; golden slide 2 should be 'Agenda'
    try:
        slide2 = prs.slides[1]  # 0-indexed
        # Search all text on slide 2 for 'Agenda'
        slide2_texts = []
        for shape in slide2.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide2_texts.append(t)
        agenda_title_found = any(t.lower() == 'agenda' for t in slide2_texts)
        if agenda_title_found:
            print(f"PASS: Component 2 -- Slide 2 has 'Agenda' title (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 -- 'Agenda' title not found on slide 2. Found texts: {slide2_texts[:5]}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: 7 agenda items with correct text (0.30 points)
    # Each correct item earns 0.30/7 points
    # GATE: Only check if agenda title was found (Component 2 passed)
    try:
        if not agenda_title_found:
            print(f"FAIL: Component 3 -- Skipped: no 'Agenda' slide at position 2")
            raise ValueError("No agenda slide")

        slide2 = prs.slides[1]
        # Collect all non-empty paragraph texts from slide 2 (excluding 'Agenda' title)
        all_para_texts = []
        for shape in slide2.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and text.lower() != 'agenda':
                        all_para_texts.append(text)

        matched_count = 0
        per_item = 0.30 / 7.0
        for expected in EXPECTED_ITEMS:
            found = any(expected.lower() in t.lower() for t in all_para_texts)
            if found:
                matched_count += 1
                total_score += per_item
            else:
                print(f"FAIL: Component 3 -- Missing agenda item: '{expected}'")

        if matched_count == 7:
            print(f"PASS: Component 3 -- All 7 agenda items found (0.30 pts)")
        else:
            print(f"PARTIAL: Component 3 -- {matched_count}/7 agenda items found ({matched_count * per_item:.3f} pts)")
            print(f"  Found texts on slide 2: {all_para_texts[:10]}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Hyperlinks to correct slides (0.25 points)
    # Each item should have a hyperlink (action=ppaction://hlinksldjump) pointing to the correct slide
    # GATE: Only check if agenda title was found (Component 2 passed)
    try:
        if not agenda_title_found:
            print(f"FAIL: Component 4 -- Skipped: no 'Agenda' slide at position 2")
            raise ValueError("No agenda slide")

        slide2 = prs.slides[1]
        # Build map: paragraph text -> hyperlink target rId
        para_hyperlinks = {}
        for shape in slide2.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and text.lower() != 'agenda':
                        for run in para.runs:
                            rPr = run._r.find(qn('a:rPr'))
                            if rPr is not None:
                                hlink = rPr.find(qn('a:hlinkClick'))
                                if hlink is not None:
                                    action = hlink.get('action', '')
                                    rId = hlink.get(qn('r:id'))
                                    if 'hlinksldjump' in action and rId:
                                        para_hyperlinks[text] = rId

        # Now check relationship targets via ZIP
        rId_to_target = {}
        with zipfile.ZipFile(file_path, 'r') as zf:
            # Find slide2 rels file
            for rels_path in ['ppt/slides/_rels/slide2.xml.rels']:
                try:
                    with zf.open(rels_path) as f:
                        root = ET.fromstring(f.read())
                        for rel in root:
                            rId = rel.get('Id')
                            target = rel.get('Target', '')
                            rId_to_target[rId] = target
                except KeyError:
                    pass

        hyperlink_correct = 0
        per_link = 0.25 / 7.0
        for idx, expected_item in enumerate(EXPECTED_ITEMS):
            target_slide_num = idx + 3  # item 0 -> slide 3, etc.
            expected_target = f'slide{target_slide_num}.xml'

            # Find matching paragraph
            matched_rId = None
            for text, rId in para_hyperlinks.items():
                if expected_item.lower() in text.lower():
                    matched_rId = rId
                    break

            if matched_rId is None:
                print(f"FAIL: Component 4 -- No hyperlink found for '{expected_item}'")
                continue

            actual_target = rId_to_target.get(matched_rId, '')
            if expected_target in actual_target:
                hyperlink_correct += 1
                total_score += per_link
            else:
                print(f"FAIL: Component 4 -- '{expected_item}' links to '{actual_target}', expected '{expected_target}'")

        if hyperlink_correct == 7:
            print(f"PASS: Component 4 -- All 7 hyperlinks point to correct slides (0.25 pts)")
        else:
            print(f"PARTIAL: Component 4 -- {hyperlink_correct}/7 correct hyperlinks ({hyperlink_correct * per_link:.3f} pts)")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Font formatting - 18pt, blue, underlined (0.15 points)
    # Check that agenda items are in 18pt (228600 EMU), blue (0000FF), and underlined
    # GATE: Only check if agenda title was found (Component 2 passed)
    try:
        if not agenda_title_found:
            print(f"FAIL: Component 5 -- Skipped: no 'Agenda' slide at position 2")
            raise ValueError("No agenda slide")

        slide2 = prs.slides[1]
        formatted_count = 0
        per_fmt = 0.15 / 7.0
        for shape in slide2.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and text.lower() != 'agenda':
                        # Check if this text matches an expected item
                        is_agenda_item = any(exp.lower() in text.lower() for exp in EXPECTED_ITEMS)
                        if not is_agenda_item:
                            continue
                        for run in para.runs:
                            if not run.text.strip():
                                continue
                            checks_pass = 0
                            total_checks = 3
                            # Check font size: 18pt = 228600 EMU
                            if run.font.size is not None and abs(run.font.size - 228600) < 1000:
                                checks_pass += 1
                            # Check underline
                            if run.font.underline:
                                checks_pass += 1
                            # Check blue color
                            try:
                                if run.font.color.type is not None:
                                    rgb_str = str(run.font.color.rgb).upper()
                                    if rgb_str == '0000FF':
                                        checks_pass += 1
                            except Exception:
                                pass

                            if checks_pass == total_checks:
                                formatted_count += 1
                            else:
                                size_val = run.font.size
                                underline_val = run.font.underline
                                try:
                                    color_val = str(run.font.color.rgb) if run.font.color.type is not None else 'N/A'
                                except:
                                    color_val = 'N/A'
                                print(f"FAIL: Component 5 -- '{text}' formatting: size={size_val}, underline={underline_val}, color={color_val} (need 228600, True, 0000FF)")
                            break  # Only check first non-empty run per paragraph

        if formatted_count >= 7:
            total_score += 0.15
            print(f"PASS: Component 5 -- All 7 items have correct formatting (0.15 pts)")
        elif formatted_count > 0:
            score_add = formatted_count * per_fmt
            total_score += score_add
            print(f"PARTIAL: Component 5 -- {formatted_count}/7 items correctly formatted ({score_add:.3f} pts)")
        else:
            print(f"FAIL: Component 5 -- No items have correct formatting")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.3f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
