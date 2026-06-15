"""
Initial Setup: Create Startup_Pitch.pptx with 9 slides (no agenda slide)
Task ID: impress_ps_042
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_042'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    """Add a title slide (layout 0)."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points):
    """Add a title+content slide (layout 1) with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(18)
    return slide


def create_initial():
    prs = Presentation()

    # Slide 1: Title Slide
    add_title_slide(prs, "NovaTech: Redefining Urban Mobility",
                    "Series A Funding Pitch  |  Q2 2025")

    # Slide 2: Problem Statement
    add_content_slide(prs, "Problem Statement", [
        "Urban commuters waste an average of 54 minutes daily in traffic",
        "Public transit systems are aging and underfunded in 73% of metro areas",
        "Ride-sharing costs have increased 42% since 2022",
        "Carbon emissions from personal vehicles remain the #1 urban pollutant",
        "Existing micro-mobility solutions lack last-mile integration",
    ])

    # Slide 3: Our Solution
    add_content_slide(prs, "Our Solution", [
        "NovaTech Smart Transit: AI-powered multimodal routing platform",
        "Seamless integration of e-bikes, e-scooters, and shuttle networks",
        "Real-time demand prediction reduces wait times by 68%",
        "Single subscription model: $49/month for unlimited rides",
        "Partnerships with 12 municipal transit agencies already signed",
    ])

    # Slide 4: Market Opportunity
    add_content_slide(prs, "Market Opportunity", [
        "Total addressable market: $847B globally by 2028",
        "Urban mobility-as-a-service growing at 19.2% CAGR",
        "Target: 25 US metro areas with 500K+ population",
        "Early traction: 38,000 active users in pilot cities (Austin, Denver)",
        "Revenue per user: $588/year with 14-month average retention",
    ])

    # Slide 5: Competitive Landscape
    add_content_slide(prs, "Competitive Landscape", [
        "Direct competitors: Lime, Bird, Via, Moovit",
        "Our edge: Only platform combining routing + fleet + payments",
        "Patent-pending demand forecasting algorithm (3 patents filed)",
        "Net Promoter Score: 72 vs industry average of 31",
        "Switching cost advantage through transit agency integrations",
    ])

    # Slide 6: Financial Projections
    add_content_slide(prs, "Financial Projections", [
        "2024 Revenue: $2.1M (actual) | 2025 Projected: $8.7M",
        "Gross margin: 62% (target 75% at scale)",
        "Customer acquisition cost: $23 (down from $41 in 2023)",
        "Burn rate: $380K/month | Runway: 18 months at current pace",
        "Path to profitability: Q3 2026 with 150K active users",
    ])

    # Slide 7: Our Team
    add_content_slide(prs, "Our Team", [
        "CEO: Dr. Priya Sharma - Former VP Engineering at Uber, MIT PhD",
        "CTO: James Rivera - Ex-Google Maps, 15 years in geospatial AI",
        "COO: Lisa Nakamura - Scaled Bird from 5 to 200 cities",
        "Head of Partnerships: David Chen - Former transit director, City of LA",
        "Advisory Board: 5 industry veterans from mobility and fintech",
    ])

    # Slide 8: The Ask
    add_content_slide(prs, "The Ask", [
        "Raising $12M Series A at $48M pre-money valuation",
        "Lead investor commitment: $5M from Velocity Ventures",
        "Use of funds: 40% engineering, 30% market expansion, 20% ops, 10% reserve",
        "Milestone targets: 100K users, 10 new cities, $5M ARR by Q4 2025",
        "Board seat offered to lead investor",
    ])

    # Slide 9: Thank You / Contact
    add_title_slide(prs, "Thank You",
                    "Dr. Priya Sharma  |  priya@novatech.io  |  novatech.io")

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
