"""
Reward Script: Delete slide 3 from the presentation
Task ID: impstruct_002
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Presentation has exactly 4 slides (down from 5)
  Component 2 (0.3): No slide contains "Outdated Info" title text
  Component 3 (0.3): Slide 3 (index 2) now has "Projects" as its title
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'impstruct_002'


def get_slide_title(slide):
    """Extract the title text from a slide, checking title shape and all text shapes."""
    # Check the title placeholder first
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        return slide.shapes.title.text.strip()
    # Fallback: first text shape
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                return text
    return ""


def get_all_slide_texts(slide):
    """Get all text content from a slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return texts


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Presentation has exactly 4 slides (0.4 points)
    # Initial has 5 slides; after deleting slide 3, should have 4
    try:
        if num_slides == 4:
            print(f"PASS: Component 1 — Slide count is 4 (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 — Expected 4 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: No slide contains "Outdated Info" text (0.3 points)
    # The deleted slide 3 had title "Outdated Info" — it must be gone
    try:
        # Collect slide indices where "Outdated Info" appears
        outdated_slides = [
            i + 1 for i, slide in enumerate(prs.slides)
            if any("Outdated Info" in t for t in get_all_slide_texts(slide))
        ]

        if len(outdated_slides) > 0:
            print(f"FAIL: Component 2 — Found 'Outdated Info' on slide(s) {outdated_slides}")
        elif num_slides < 5:
            # Only award points if slide count also changed (prevents false pass on unmodified file
            # that somehow lacks the text)
            print(f"PASS: Component 2 — 'Outdated Info' slide not present (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 — 'Outdated Info' text not found but slide count is still {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 3 (index 2) is now "Projects" (0.3 points)
    # After deleting old slide 3 ("Outdated Info"), "Projects" shifts from position 4 to position 3
    try:
        if num_slides >= 3:
            slide3_title = get_slide_title(prs.slides[2])
            if slide3_title == "Projects":
                print(f"PASS: Component 3 — Slide 3 title is 'Projects' (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 3 — Expected slide 3 title 'Projects', found '{slide3_title}'")
        else:
            print(f"FAIL: Component 3 — Not enough slides to check slide 3 (only {num_slides})")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
