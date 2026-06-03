"""
Initial Setup: Create a 5-slide team introduction presentation
Task ID: impstruct_002
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impstruct_002'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Welcome ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    slide1.shapes.title.text = "Welcome"
    slide1.placeholders[1].text = "Quarterly Team Introduction\nApril 2025"

    # --- Slide 2: Team Lead ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide2.shapes.title.text = "Team Lead"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Sarah Chen — Director of Engineering"
    p2 = body2.add_paragraph()
    p2.text = "15 years of experience in distributed systems"
    p2.level = 1
    p3 = body2.add_paragraph()
    p3.text = "Previously at Stripe, Google, and Databricks"
    p3.level = 1
    p4 = body2.add_paragraph()
    p4.text = "Joined Acme Corp in January 2024"
    p4.level = 1

    # --- Slide 3: Outdated Info ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide3.shapes.title.text = "Outdated Info"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Previous Office Location: 450 Market Street, Suite 200"
    p = body3.add_paragraph()
    p.text = "Old phone system extension chart (decommissioned)"
    p.level = 1
    p = body3.add_paragraph()
    p.text = "Legacy project codenames no longer in use"
    p.level = 1
    p = body3.add_paragraph()
    p.text = "This slide contains outdated information and should be removed."
    p.level = 0

    # --- Slide 4: Projects ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide4.shapes.title.text = "Projects"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Current Active Projects"
    p = body4.add_paragraph()
    p.text = "Project Atlas — Cloud migration (Q2 2025)"
    p.level = 1
    p = body4.add_paragraph()
    p.text = "Project Beacon — Customer analytics dashboard"
    p.level = 1
    p = body4.add_paragraph()
    p.text = "Project Cipher — Security audit and compliance"
    p.level = 1
    p = body4.add_paragraph()
    p.text = "Project Delta — Mobile app redesign"
    p.level = 1

    # --- Slide 5: Contact ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide5.shapes.title.text = "Contact"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Reach Us"
    p = body5.add_paragraph()
    p.text = "Email: team-engineering@acmecorp.com"
    p.level = 1
    p = body5.add_paragraph()
    p.text = "Slack: #engineering-general"
    p.level = 1
    p = body5.add_paragraph()
    p.text = "Office: Building C, Floor 3"
    p.level = 1
    p = body5.add_paragraph()
    p.text = "Office Hours: Mon-Fri 9 AM - 6 PM PST"
    p.level = 1

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Number of slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
