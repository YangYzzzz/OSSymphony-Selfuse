"""
Reward Script: Insert pie chart on slide 4 with market share data
Task ID: impress_ps_045
Domain: libreoffice_impress
Scoring:
  Component 1: Pie chart exists on slide 4 (0.20 pts)
  Component 2: Correct 4 categories and percentage values (0.25 pts)
  Component 3: Acme Corp slice is exploded (0.20 pts)
  Component 4: Percentage labels on slices (0.20 pts)
  Component 5: Legend is present (0.15 pts)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_045'


def persist_app_state(domain):
    """Save any unsaved LibreOffice state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print("PERSIST: ctrl+s sent for %s" % domain)
        except Exception as e:
            print("PERSIST_WARN: save hook failed: %s" % e)


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation
    from pptx.oxml.ns import qn

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print("CRITICAL: Cannot load file %s: %s" % (file_path, e))
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 4 slides
    if len(prs.slides) < 4:
        print("FAIL: Presentation has fewer than 4 slides (%d found)" % len(prs.slides))
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # slide 4 (0-indexed)

    # Find chart shape on slide 4
    chart_shape = None
    for shape in slide.shapes:
        if shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Pie chart exists on slide 4 (0.20 points)
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            # Check it's a pie chart (type 5 = PIE)
            chart_type_val = chart.chart_type
            # PIE is 5, EXPLODED_PIE is -4102 or similar
            is_pie = 'PIE' in str(chart_type_val)
            if is_pie:
                print("PASS: Component 1 — Pie chart found on slide 4 (0.20 pts)")
                total_score += 0.20
            else:
                print("FAIL: Component 1 — Chart on slide 4 is not a pie chart, type=%s" % chart_type_val)
        else:
            print("FAIL: Component 1 — No chart found on slide 4")
    except Exception as e:
        print("ERROR: Component 1 — %s" % e)

    if chart_shape is None:
        print("\nScore: %.2f/1.0" % total_score)
        print("REWARD: %.1f" % total_score)
        return total_score

    chart = chart_shape.chart

    # Component 2: Correct categories and values (0.25 points)
    # Expected: Acme Corp 35%, Competitor A 25%, Competitor B 20%, Others 20%
    try:
        plot = chart.plots[0]
        series = plot.series[0]
        values = list(series.values)

        try:
            categories = [str(c) for c in plot.categories]
        except Exception:
            categories = []

        expected_cats = ['Acme Corp', 'Competitor A', 'Competitor B', 'Others']
        expected_vals = [35.0, 25.0, 20.0, 20.0]

        # Check categories match (case-insensitive, stripped)
        cats_match = len(categories) == 4
        if cats_match:
            for exp_c, act_c in zip(expected_cats, categories):
                if exp_c.lower().strip() != act_c.lower().strip():
                    cats_match = False
                    break

        # Check values match (allow small tolerance)
        vals_match = len(values) == 4
        if vals_match:
            for exp_v, act_v in zip(expected_vals, values):
                if abs(float(exp_v) - float(act_v)) > 0.5:
                    vals_match = False
                    break

        if cats_match and vals_match:
            total_score += 0.25
            print("PASS: Component 2 — Correct 4 categories and values (0.25 pts)")
            print("  Categories: %s" % categories)
            print("  Values: %s" % values)
        elif vals_match:
            total_score += 0.10
            print("PARTIAL: Component 2 — Values correct but categories mismatch (0.10 pts)")
            print("  Expected categories: %s" % expected_cats)
            print("  Found categories: %s" % categories)
            print("  Values: %s" % values)
        else:
            print("FAIL: Component 2 — Categories or values incorrect")
            print("  Expected categories: %s" % expected_cats)
            print("  Found categories: %s" % categories)
            print("  Expected values: %s" % expected_vals)
            print("  Found values: %s" % values)
    except Exception as e:
        print("ERROR: Component 2 — %s" % e)

    # Component 3: Acme Corp slice (index 0) is exploded (0.20 points)
    try:
        plot = chart.plots[0]
        series = plot.series[0]
        ser_elem = series._element

        # Check per-point explosion via XML
        dpts = ser_elem.findall(qn('c:dPt'))
        acme_explosion_val = 0
        for dpt in dpts:
            idx_elem = dpt.find(qn('c:idx'))
            exp_elem = dpt.find(qn('c:explosion'))
            if idx_elem is not None and idx_elem.get('val') == '0' and exp_elem is not None:
                acme_explosion_val = int(exp_elem.get('val', '0'))

        # Also check series-level explosion as fallback
        if acme_explosion_val == 0:
            ser_exp = ser_elem.find(qn('c:explosion'))
            if ser_exp is not None:
                acme_explosion_val = int(ser_exp.get('val', '0'))

        if acme_explosion_val > 0:
            print("PASS: Component 3 — Acme Corp slice exploded (explosion=%d) (0.20 pts)" % acme_explosion_val)
            total_score += 0.20
        else:
            print("FAIL: Component 3 — Acme Corp slice is not exploded")
    except Exception as e:
        print("ERROR: Component 3 — %s" % e)

    # Component 4: Percentage labels displayed on slices (0.20 points)
    try:
        plot = chart.plots[0]
        dl = plot.data_labels
        has_pct_labels = dl.show_percentage
        if has_pct_labels:
            print("PASS: Component 4 — Percentage labels enabled on slices (0.20 pts)")
            total_score += 0.20
        else:
            # Check via XML as fallback
            plot_elem = plot._element
            dLbls = plot_elem.find(qn('c:dLbls'))
            if dLbls is not None:
                show_pct = dLbls.find(qn('c:showPercent'))
                if show_pct is not None and show_pct.get('val') == '1':
                    print("PASS: Component 4 — Percentage labels found via XML (0.20 pts)")
                    total_score += 0.20
                else:
                    print("FAIL: Component 4 — showPercent not enabled (val=%s)" % (
                        show_pct.get('val') if show_pct is not None else 'element missing'))
            else:
                print("FAIL: Component 4 — No data labels element found")
    except Exception as e:
        print("ERROR: Component 4 — %s" % e)

    # Component 5: Legend is present (0.15 points)
    try:
        if chart.has_legend:
            print("PASS: Component 5 — Chart legend is present (0.15 pts)")
            total_score += 0.15
        else:
            print("FAIL: Component 5 — Chart has no legend")
    except Exception as e:
        print("ERROR: Component 5 — %s" % e)

    final_score = min(round(total_score, 2), 1.0)
    print("\nScore: %.2f/1.0" % final_score)
    print("REWARD: %.1f" % final_score)
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = '%s/%s.pptx' % (WORKDIR, TASK_ID)
if not os.path.exists(file_path):
    print("File not found: %s" % file_path)
    print("REWARD: 0.0")
else:
    verify_task(file_path)
