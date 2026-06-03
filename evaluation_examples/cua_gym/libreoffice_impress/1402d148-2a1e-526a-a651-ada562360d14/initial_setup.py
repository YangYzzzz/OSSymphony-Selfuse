"""
Initial Setup: Create Annual_Review_2025.pptx with 12 slides, slide 4 titled 'Market Share' with empty content area.
Task ID: impress_ps_045
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_045'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Annual Review 2025"
    slide1.placeholders[1].text = "Acme Corporation\nQ1-Q4 Performance Summary"

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Acme Corporation achieved record growth in 2025, with total revenue reaching $2.4 billion."
    p = body2.add_paragraph()
    p.text = "Key highlights include expansion into 3 new markets, a 15% increase in customer retention, and the successful launch of our flagship product line."
    p = body2.add_paragraph()
    p.text = "Employee satisfaction scores rose to 87%, reflecting our investment in workplace culture."

    # --- Slide 3: Revenue Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Revenue Overview"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Q1: $520M (+8% YoY)"
    for item in ["Q2: $610M (+12% YoY)", "Q3: $640M (+14% YoY)", "Q4: $630M (+11% YoY)", "Full Year: $2.4B (+11.3% YoY)"]:
        p = body3.add_paragraph()
        p.text = item

    # --- Slide 4: Market Share (EMPTY - NO CHART) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title text box at top
    title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Market Share"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x3B, 0x4E)

    # --- Slide 5: Product Performance ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Product Performance"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "CloudSync Pro: $890M revenue, 2.1M active users"
    for item in ["DataVault Enterprise: $620M revenue, 850K licenses", "SmartFlow Analytics: $410M revenue, 1.4M monthly queries", "SecureNet Platform: $480M revenue, 3,200 enterprise clients"]:
        p = body5.add_paragraph()
        p.text = item

    # --- Slide 6: Customer Acquisition ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Customer Acquisition"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "New customers acquired: 12,400 (+18% vs 2024)"
    for item in ["Enterprise segment grew 22%, adding 340 Fortune 500 accounts", "SMB segment expanded 15% through partner channel programs", "Customer acquisition cost decreased 8% through digital marketing optimization"]:
        p = body6.add_paragraph()
        p.text = item

    # --- Slide 7: Regional Breakdown ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Regional Breakdown"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "North America: $1.2B (50% of total revenue)"
    for item in ["Europe: $720M (30%)", "Asia Pacific: $360M (15%)", "Rest of World: $120M (5%)"]:
        p = body7.add_paragraph()
        p.text = item

    # --- Slide 8: R&D Investments ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "R&D Investments"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Total R&D spend: $380M (15.8% of revenue)"
    for item in ["AI/ML capabilities: $120M allocated to next-gen features", "Cloud infrastructure modernization: $95M", "Security enhancements: $85M", "Patent filings: 47 new patents granted in 2025"]:
        p = body8.add_paragraph()
        p.text = item

    # --- Slide 9: Employee Metrics ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Employee Metrics"
    body9 = slide9.placeholders[1].text_frame
    body9.text = "Total headcount: 8,450 (+620 net new hires)"
    for item in ["Engineering: 3,200 (38%)", "Sales & Marketing: 2,100 (25%)", "Operations: 1,850 (22%)", "G&A: 1,300 (15%)", "Voluntary turnover: 8.2% (industry avg: 13.5%)"]:
        p = body9.add_paragraph()
        p.text = item

    # --- Slide 10: Strategic Partnerships ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "Strategic Partnerships"
    body10 = slide10.placeholders[1].text_frame
    body10.text = "Microsoft Azure: Expanded integration for CloudSync Pro"
    for item in ["Salesforce: Joint go-to-market program launched Q2", "Deloitte: Implementation partner for enterprise deployments", "AWS: Marketplace listing drove $45M in new pipeline"]:
        p = body10.add_paragraph()
        p.text = item

    # --- Slide 11: 2026 Outlook ---
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    slide11.shapes.title.text = "2026 Outlook"
    body11 = slide11.placeholders[1].text_frame
    body11.text = "Revenue target: $2.8B (+17% growth)"
    for item in ["Launch of 2 new product lines in Q1 and Q3", "Expansion into Latin America and Middle East markets", "Planned acquisition pipeline: $200M allocated", "IPO preparation activities to begin Q2 2026"]:
        p = body11.add_paragraph()
        p.text = item

    # --- Slide 12: Thank You ---
    slide12 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    thank_box = slide12.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = thank_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x3B, 0x4E)
    p2 = tf.add_paragraph()
    p2.text = "Questions & Discussion"
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.name = "Calibri"
    run2.font.size = Pt(24)
    run2.font.color.rgb = RGBColor(0x6B, 0x7B, 0x8D)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
