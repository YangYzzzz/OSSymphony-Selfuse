"""
Initial Setup: Philosophy 101 presentation with 6 slides. Slide 3 has title and content but no visual separator.
Task ID: impress_teach_018
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_018'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Philosophy 101"
    slide1.placeholders[1].text = "An Introduction to Western Philosophical Thought\nFall 2025 — Prof. Elena Vasquez"

    # --- Slide 2: What is Philosophy? ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "What is Philosophy?"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Philosophy — from the Greek philosophia, meaning 'love of wisdom'"
    p2a = tf2.add_paragraph()
    p2a.text = "The systematic study of fundamental questions about existence, knowledge, values, reason, and language"
    p2a.level = 1
    p2b = tf2.add_paragraph()
    p2b.text = "Unlike empirical sciences, philosophy relies primarily on rational argumentation and critical analysis"
    p2b.level = 1
    p2c = tf2.add_paragraph()
    p2c.text = "Philosophical inquiry dates back to at least the 6th century BCE in ancient Greece, India, and China"
    p2c.level = 1

    # --- Slide 3: Major Branches (NO line separator here) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Major Branches of Philosophy"
    tf3 = slide3.placeholders[1].text_frame
    branches = [
        ("Metaphysics", "The study of the fundamental nature of reality, including the relationship between mind and matter"),
        ("Epistemology", "The theory of knowledge — its nature, scope, and the justification of belief"),
        ("Ethics", "The branch concerned with moral principles that govern behavior and decision-making"),
        ("Logic", "The study of valid reasoning, argument structure, and the principles of correct inference"),
        ("Aesthetics", "The philosophical study of beauty, art, and the nature of taste and perception"),
    ]
    tf3.text = f"{branches[0][0]}: {branches[0][1]}"
    for name, desc in branches[1:]:
        p = tf3.add_paragraph()
        p.text = f"{name}: {desc}"
        p.level = 0

    # --- Slide 4: Key Thinkers ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Key Thinkers in Western Philosophy"
    tf4 = slide4.placeholders[1].text_frame
    thinkers = [
        "Socrates (470–399 BCE) — The Socratic method of questioning assumptions",
        "Plato (428–348 BCE) — Theory of Forms and The Republic",
        "Aristotle (384–322 BCE) — Formal logic, virtue ethics, and empirical inquiry",
        "Immanuel Kant (1724–1804) — Critique of Pure Reason and the categorical imperative",
        "Simone de Beauvoir (1908–1986) — Existentialist feminism and The Second Sex",
        "Ludwig Wittgenstein (1889–1951) — Language games and the limits of thought",
    ]
    tf4.text = thinkers[0]
    for t in thinkers[1:]:
        p = tf4.add_paragraph()
        p.text = t
        p.level = 0

    # --- Slide 5: Philosophical Methods ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Philosophical Methods"
    tf5 = slide5.placeholders[1].text_frame
    methods = [
        "Dialectic — structured dialogue between opposing viewpoints to reveal contradictions",
        "Thought Experiments — imagined scenarios that isolate key philosophical intuitions",
        "Conceptual Analysis — breaking down complex ideas into constituent elements",
        "Phenomenological Description — careful first-person accounts of conscious experience",
        "Formal Logic — symbolic representation and evaluation of argument validity",
    ]
    tf5.text = methods[0]
    for m in methods[1:]:
        p = tf5.add_paragraph()
        p.text = m
        p.level = 0

    # --- Slide 6: Course Overview ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Course Overview & Assessment"
    tf6 = slide6.placeholders[1].text_frame
    items = [
        "Weeks 1–4: Ancient Philosophy — Pre-Socratics through Aristotle",
        "Weeks 5–8: Medieval & Early Modern — Augustine, Aquinas, Descartes, Hume",
        "Weeks 9–12: Modern & Contemporary — Kant, Hegel, existentialism, analytic philosophy",
        "Assessment: Two essays (30% each), class participation (15%), final exam (25%)",
        "Office hours: Tuesdays & Thursdays, 2:00–3:30 PM, Room 412 Humanities Building",
    ]
    tf6.text = items[0]
    for item in items[1:]:
        p = tf6.add_paragraph()
        p.text = item
        p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
