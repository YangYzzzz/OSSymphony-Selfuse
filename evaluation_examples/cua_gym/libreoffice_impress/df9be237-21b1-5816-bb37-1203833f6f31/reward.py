"""
Reward Script: Insert a horizontal line shape on slide 3 below the title
Task ID: impress_teach_018
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3) - A line shape exists on slide 3
  Component 2 (0.3) - Line is horizontal and positioned below title area
  Component 3 (0.2) - Line thickness is 2pt (25400 EMU)
  Component 4 (0.2) - Line color is dark red #8B0000
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_018'


def persist_app_state(domain: str):
    """Save any unsaved GUI edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Pt, Emu
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"PRECONDITION FAIL: Need at least 3 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[2]  # slide 3 (0-indexed)

    # Find all line shapes on slide 3
    line_shapes = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            line_shapes.append(shape)

    # Component 1: A line shape exists on slide 3 (0.3 points)
    try:
        if len(line_shapes) > 0:
            print(f"PASS: Component 1 -- Line shape found on slide 3 ({len(line_shapes)} line(s)) (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 -- No line shape found on slide 3")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    if len(line_shapes) == 0:
        # No line found, remaining components fail
        print(f"FAIL: Component 2 -- No line to check position")
        print(f"FAIL: Component 3 -- No line to check thickness")
        print(f"FAIL: Component 4 -- No line to check color")
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Use the first line shape found (or find the best candidate)
    # Pick the line that best matches: horizontal, below title area
    best_line = None
    for ls in line_shapes:
        # A horizontal line has height == 0 or very small relative to width
        if ls.height <= ls.width * 0.05:  # essentially horizontal
            best_line = ls
            break
    if best_line is None:
        best_line = line_shapes[0]  # fallback to first line

    # Component 2: Line is horizontal and positioned below title area (0.3 points)
    # Title area bottom: title top (274638) + title height (1143000) = 1417638 EMU
    # Line should be below title (top > ~1400000) and roughly horizontal (height ~= 0)
    try:
        is_horizontal = best_line.height <= best_line.width * 0.05
        title_bottom = 274638 + 1143000  # ~1417638 EMU
        # Allow some tolerance - line should be below title and above content
        below_title = best_line.top >= title_bottom - 200000  # some tolerance
        above_content_midpoint = best_line.top <= 3000000  # well above slide midpoint

        # Also check it spans a reasonable width (at least 50% of slide width)
        slide_width = prs.slide_width  # typically 9144000 EMU (10 inches)
        spans_width = best_line.width >= slide_width * 0.5

        if is_horizontal and below_title and spans_width:
            print(f"PASS: Component 2 -- Line is horizontal (h={best_line.height}), "
                  f"below title (top={best_line.top}), spans width ({best_line.width}) (0.3 pts)")
            total_score += 0.3
        else:
            reasons = []
            if not is_horizontal:
                reasons.append(f"not horizontal (height={best_line.height}, width={best_line.width})")
            if not below_title:
                reasons.append(f"not below title (top={best_line.top}, title_bottom={title_bottom})")
            if not spans_width:
                reasons.append(f"too narrow (width={best_line.width}, need >= {slide_width * 0.5})")
            print(f"FAIL: Component 2 -- {'; '.join(reasons)}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Line thickness is 2pt (25400 EMU) (0.2 points)
    try:
        line_width = best_line.line.width
        # 2pt = 25400 EMU. Allow some tolerance (1.5pt to 2.5pt range)
        expected_width = 25400  # 2pt
        tolerance = 6350  # 0.5pt
        if line_width is not None and abs(line_width - expected_width) <= tolerance:
            print(f"PASS: Component 3 -- Line thickness is {line_width} EMU "
                  f"({line_width/12700:.1f}pt), expected 2pt (0.2 pts)")
            total_score += 0.2
        else:
            actual_pt = line_width / 12700 if line_width else 0
            print(f"FAIL: Component 3 -- Line thickness is {line_width} EMU ({actual_pt:.1f}pt), expected 2pt (25400 EMU)")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Line color is dark red #8B0000 (0.2 points)
    try:
        line_color = None
        if best_line.line.color and best_line.line.color.type is not None:
            line_color = str(best_line.line.color.rgb).upper()

        expected_color = "8B0000"
        if line_color is not None and line_color == expected_color:
            print(f"PASS: Component 4 -- Line color is #{line_color} (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 4 -- Line color is {line_color}, expected #{expected_color}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
