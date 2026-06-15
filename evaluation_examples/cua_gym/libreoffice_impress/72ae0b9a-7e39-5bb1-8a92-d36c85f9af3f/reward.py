"""
Reward Script: Portfolio Teaser Presentation
Task ID: impress_wf_006
Domain: libreoffice_impress
Scoring:
  Component 1 (0.15): File exists on Desktop with exactly 3 slides
  Component 2 (0.25): Slide 1 has black background + white centered 'My Portfolio' title
  Component 3 (0.25): Slide 2 has 4 images in 2x2 grid with white border rectangles
  Component 4 (0.15): Slide 3 has centered contact info (name, email, phone)
  Component 5 (0.20): All 3 slides have Push Left transitions
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_006'
FILE_PATH = os.path.join(WORKDIR, 'Desktop', 'Portfolio_Teaser.pptx')


def persist_app_state(domain: str):
    """Best-effort save in case file is still open in LibreOffice."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: File has exactly 3 slides (0.15 points)
    # Initial env has NO file on Desktop, so this differentiates.
    try:
        num_slides = len(prs.slides)
        if num_slides == 3:
            print(f"PASS: Component 1 — File has 3 slides (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — Expected 3 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 1 has black background + white centered 'My Portfolio' (0.25 points)
    try:
        slide1 = prs.slides[0]
        comp2_score = 0.0

        # Check black background
        fill = slide1.background.fill
        bg_black = False
        if fill.type == 1:  # SOLID
            bg_rgb = str(fill.fore_color.rgb)
            if bg_rgb == '000000':
                bg_black = True

        # Check 'My Portfolio' text, white, centered
        title_found = False
        title_white = False
        title_centered = False
        for shape in slide1.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if 'my portfolio' in text.lower():
                        title_found = True
                        # Check alignment
                        if para.alignment is not None and para.alignment == 2:  # CENTER
                            title_centered = True
                        # Check white color on runs
                        for run in para.runs:
                            try:
                                if run.font.color.type is not None:
                                    rgb = str(run.font.color.rgb)
                                    if rgb == 'FFFFFF':
                                        title_white = True
                            except:
                                pass

        if bg_black:
            comp2_score += 0.10
        else:
            print(f"FAIL: Component 2a — Slide 1 background is not black")

        if title_found and title_white and title_centered:
            comp2_score += 0.15
        else:
            print(f"FAIL: Component 2b — title_found={title_found}, white={title_white}, centered={title_centered}")

        if comp2_score > 0:
            print(f"PASS: Component 2 — Slide 1 black bg + white centered title ({comp2_score} pts)")
            total_score += comp2_score
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 2 has 4 images with white border rectangles in 2x2 grid (0.25 points)
    try:
        slide2 = prs.slides[1]
        comp3_score = 0.0

        # Count images
        images = [s for s in slide2.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        num_images = len(images)

        # Count white-filled rectangles (borders)
        white_rects = []
        for s in slide2.shapes:
            if s.shape_type == 1:  # AUTO_SHAPE (rectangle)
                try:
                    if s.fill.type == 1:  # SOLID
                        rgb = str(s.fill.fore_color.rgb)
                        if rgb == 'FFFFFF':
                            white_rects.append(s)
                except:
                    pass

        # Check 4 images exist
        if num_images == 4:
            comp3_score += 0.10
        else:
            print(f"FAIL: Component 3a — Expected 4 images, found {num_images}")

        # Check 4 white border rectangles
        if len(white_rects) == 4:
            comp3_score += 0.05
        else:
            print(f"FAIL: Component 3b — Expected 4 white border rects, found {len(white_rects)}")

        # Check 2x2 grid arrangement: images should form 2 rows x 2 cols
        # Verify by checking that there are 2 distinct top positions and 2 distinct left positions
        if num_images == 4:
            tops = sorted(set(img.top for img in images))
            lefts = sorted(set(img.left for img in images))
            # Allow some tolerance — group tops and lefts within 5% of slide height/width
            def cluster(values, tolerance):
                clusters = []
                for v in sorted(values):
                    if clusters and abs(v - clusters[-1][-1]) < tolerance:
                        clusters[-1].append(v)
                    else:
                        clusters.append([v])
                return len(clusters)

            top_clusters = cluster([img.top for img in images], prs.slide_height * 0.05)
            left_clusters = cluster([img.left for img in images], prs.slide_width * 0.05)

            if top_clusters == 2 and left_clusters == 2:
                comp3_score += 0.10
            else:
                print(f"FAIL: Component 3c — Grid check: {top_clusters} rows, {left_clusters} cols (expected 2x2)")

        if comp3_score > 0:
            print(f"PASS: Component 3 — Slide 2 image grid ({comp3_score} pts)")
            total_score += comp3_score
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 3 has centered contact info (name, email, phone) (0.15 points)
    try:
        slide3 = prs.slides[2]
        comp4_score = 0.0

        # Collect all text from the slide
        all_text = []
        has_centered_textbox = False
        has_email = False
        has_phone = False
        has_name = False

        for shape in slide3.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        all_text.append(text)
                        # Check centering
                        if para.alignment is not None and para.alignment == 2:  # CENTER
                            has_centered_textbox = True
                        # Check for email pattern
                        if '@' in text and '.' in text:
                            has_email = True
                        # Check for phone pattern (digits, parentheses, dashes, plus)
                        import re
                        if re.search(r'[\d\(\)\-\+\s]{7,}', text):
                            has_phone = True
                        # Any non-email, non-phone text is likely the name
                        if '@' not in text and not re.search(r'^\+?\d[\d\(\)\-\s]{6,}$', text.strip()):
                            has_name = True

        if has_name and has_email and has_phone:
            comp4_score += 0.10
        else:
            print(f"FAIL: Component 4a — name={has_name}, email={has_email}, phone={has_phone}")

        if has_centered_textbox:
            comp4_score += 0.05
        else:
            print(f"FAIL: Component 4b — Contact text not centered")

        if comp4_score > 0:
            print(f"PASS: Component 4 — Slide 3 contact info ({comp4_score} pts)")
            total_score += comp4_score
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: All 3 slides have Push Left transitions (0.20 points)
    try:
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        push_count = 0
        with zipfile.ZipFile(file_path, 'r') as zf:
            for si in range(1, 4):
                try:
                    with zf.open(f'ppt/slides/slide{si}.xml') as f:
                        root = ET.parse(f).getroot()
                        tr = root.find('.//p:transition', ns)
                        if tr is not None:
                            push = tr.find('p:push', ns)
                            if push is not None:
                                direction = push.attrib.get('dir', '')
                                if direction == 'l':
                                    push_count += 1
                                    print(f"  Slide {si}: Push Left transition found")
                                else:
                                    print(f"  Slide {si}: Push transition found but dir={direction}")
                            else:
                                print(f"  Slide {si}: Transition exists but not Push")
                        else:
                            print(f"  Slide {si}: No transition found")
                except Exception as e:
                    print(f"  Slide {si}: Error checking transition: {e}")

        if push_count == 3:
            print(f"PASS: Component 5 — All 3 slides have Push Left transitions (0.20 pts)")
            total_score += 0.20
        elif push_count > 0:
            partial = round(0.20 * push_count / 3, 2)
            print(f"PARTIAL: Component 5 — {push_count}/3 slides have Push Left transitions ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 — No Push Left transitions found")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(FILE_PATH)
