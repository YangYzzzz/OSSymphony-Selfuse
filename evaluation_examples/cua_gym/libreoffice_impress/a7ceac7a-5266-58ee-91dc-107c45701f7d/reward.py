"""
Reward Script: Update the title of slide 2 to 'Q3 Financial Results' and justify the text.
Task ID: osworld_impress_title_set_aligned_006
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 2 title text is 'Q3 Financial Results' (0.6 points)
  Component 2: Slide 2 title alignment is JUSTIFY (0.4 points)
  Total: 1.0
"""

import os
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_set_aligned_006'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.

    Task: Update slide 2 title text to 'Q3 Financial Results' and set alignment to JUSTIFY.
    Initial state: slide 2 title is 'Q3 Results' with LEFT alignment.
    Golden state:  slide 2 title is 'Q3 Financial Results' with JUSTIFY alignment.

    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load presentation — failure here means 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition gate: presentation must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"CRITICAL: Expected at least 2 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]  # 0-indexed: slide 2

    # Find the title shape on slide 2
    title_shape = None
    for shape in slide2.shapes:
        if shape.has_text_frame and shape.name == 'Title 1':
            title_shape = shape
            break

    if title_shape is None:
        print("CRITICAL: Title shape not found on slide 2")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide 2 title text is 'Q3 Financial Results' (0.6 points)
    # Initial state has 'Q3 Results' — this FAILS on initial, PASSES on golden
    try:
        title_para = title_shape.text_frame.paragraphs[0]
        actual_title = title_para.text.strip()
        expected_title = 'Q3 Financial Results'
        if actual_title == expected_title:
            print(f"PASS: Component 1 — Slide 2 title text is '{actual_title}' (0.6 pts)")
            total_score += 0.6
        else:
            print(f"FAIL: Component 1 — Expected title '{expected_title}', found '{actual_title}'")
    except Exception as e:
        print(f"ERROR: Component 1 — Could not read slide 2 title text: {e}")

    # Component 2: Slide 2 title paragraph alignment is JUSTIFY (0.4 points)
    # Initial state has LEFT alignment (PP_ALIGN.LEFT == 1) — this FAILS on initial, PASSES on golden
    try:
        title_para = title_shape.text_frame.paragraphs[0]
        actual_alignment = title_para.alignment
        # JUSTIFY is PP_ALIGN.JUSTIFY (value 4)
        if actual_alignment == PP_ALIGN.JUSTIFY:
            print(f"PASS: Component 2 — Slide 2 title alignment is JUSTIFY ({actual_alignment}) (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 2 — Expected JUSTIFY alignment ({PP_ALIGN.JUSTIFY}), found '{actual_alignment}'")
    except Exception as e:
        print(f"ERROR: Component 2 — Could not read slide 2 title alignment: {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in VM env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
