"""
Initial Setup: 5-slide financial report presentation
Task ID: osworld_impress_title_set_aligned_006
Domain: libreoffice_impress

Creates a presentation where slide 2 has title 'Q3 Results' with LEFT alignment.
The agent task is to change the title to 'Q3 Financial Results' and justify it.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_set_aligned_006'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_para_text(tf, text, alignment=None, font_size=None, bold=None, color=None):
    """Helper to set text frame paragraph with formatting."""
    tf.clear()
    para = tf.paragraphs[0]
    if alignment is not None:
        para.alignment = alignment
    run = para.add_run()
    run.text = text
    if font_size is not None:
        run.font.size = Pt(font_size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def create_initial():
    prs = Presentation()
    # Use standard widescreen dimensions (10 x 7.5 inches)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title slide ─────────────────────────────────────────────────
    layout_title = prs.slide_layouts[0]  # Title Slide
    slide1 = prs.slides.add_slide(layout_title)
    slide1.shapes.title.text = 'Annual Financial Report 2024'
    slide1.placeholders[1].text = 'Prepared by the Finance Department\nFiscal Year Overview'

    # ── Slide 2: Q3 Results (title LEFT-aligned, NOT justified) ──────────────
    layout_content = prs.slide_layouts[1]  # Title and Content
    slide2 = prs.slides.add_slide(layout_content)

    # Set title: 'Q3 Results' with explicit LEFT alignment
    title2 = slide2.shapes.title
    title2.text = ''
    para2 = title2.text_frame.paragraphs[0]
    para2.alignment = PP_ALIGN.LEFT
    run2 = para2.add_run()
    run2.text = 'Q3 Results'
    run2.font.size = Pt(32)
    run2.font.bold = True

    # Content body
    body2 = slide2.placeholders[1]
    tf2 = body2.text_frame
    tf2.clear()
    lines2 = [
        'Revenue: $4.2M (up 12% YoY)',
        'Operating Expenses: $2.8M',
        'Net Income: $1.4M',
        'EBITDA Margin: 33%',
        'Customer Acquisition Cost: $145',
        'Recurring Revenue Rate: 87%',
    ]
    for i, line in enumerate(lines2):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = line
        p.level = 0

    # ── Slide 3: Revenue Breakdown ───────────────────────────────────────────
    slide3 = prs.slides.add_slide(layout_content)
    slide3.shapes.title.text = 'Revenue Breakdown by Region'
    body3 = slide3.placeholders[1]
    tf3 = body3.text_frame
    tf3.clear()
    lines3 = [
        'North America: $1.85M (44%)',
        'Europe: $1.22M (29%)',
        'Asia-Pacific: $0.88M (21%)',
        'Latin America: $0.25M (6%)',
        'Total: $4.20M',
    ]
    for i, line in enumerate(lines3):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = line

    # ── Slide 4: Cost Analysis ───────────────────────────────────────────────
    slide4 = prs.slides.add_slide(layout_content)
    slide4.shapes.title.text = 'Operating Cost Analysis'
    body4 = slide4.placeholders[1]
    tf4 = body4.text_frame
    tf4.clear()
    lines4 = [
        'Personnel Costs: $1.65M (59%)',
        'Infrastructure & IT: $0.42M (15%)',
        'Marketing & Sales: $0.39M (14%)',
        'Research & Development: $0.22M (8%)',
        'General & Administrative: $0.12M (4%)',
        'Cost Reduction Target Q4: 5%',
    ]
    for i, line in enumerate(lines4):
        if i == 0:
            p = tf4.paragraphs[0]
        else:
            p = tf4.add_paragraph()
        p.text = line

    # ── Slide 5: Outlook & Summary ───────────────────────────────────────────
    slide5 = prs.slides.add_slide(layout_content)
    slide5.shapes.title.text = 'Q4 Outlook & Strategic Priorities'
    body5 = slide5.placeholders[1]
    tf5 = body5.text_frame
    tf5.clear()
    lines5 = [
        'Revenue Target: $4.8M (+14% growth)',
        'New Product Launch: CloudSync Pro v2.0',
        'Expand into Southeast Asia market',
        'Headcount: +12 engineers, +5 sales',
        'Infrastructure upgrade investment: $0.35M',
        'Projected EBITDA Margin: 35%',
    ]
    for i, line in enumerate(lines5):
        if i == 0:
            p = tf5.paragraphs[0]
        else:
            p = tf5.add_paragraph()
        p.text = line

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open the file in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
