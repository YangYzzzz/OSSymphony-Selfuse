"""
FINAL REWARD SCRIPT - SUCCESS
Task: Quick question: on slide 91 of my LibreOffice Impress deck, I need to insert a table that’s exactly 2 rows by 10 columns and make sure the text in column 2 is right-aligned. What’s the fastest way to set that up?
Generated: 2025-09-10 23:22:02
Status: success
Model: azure-o3
Total Steps: 1
"""

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
import os

"""
Reward Script for LibreOffice Impress Task
Task: Verify that on slide 91 of the given PPTX deck there is a table that:
  • Has exactly 2 rows and 10 columns
  • Has text in column 2 (index 1) right-aligned in every row
The script awards a progressive score and returns exactly 1.0 only when every requirement is met.
"""

# ------------------------------------------------------------------
# Helper Function: verify right-alignment in column 2 of the table
# ------------------------------------------------------------------

def verify_table_alignment(table):
    """Return True if every paragraph in column 2 of all rows is right-aligned."""
    for r in range(len(table.rows)):
        cell = table.cell(r, 1)  # column 2 (0-based index)
        tf = cell.text_frame
        if tf is None:
            print(f"    ✗ Row {r+1}, column 2 has no text frame")
            return False
        for p in tf.paragraphs:
            if p.alignment != PP_ALIGN.RIGHT:
                print(
                    f"    ✗ Row {r+1}, column 2 paragraph not right-aligned (alignment={p.alignment})"
                )
                return False
    return True

# ------------------------------------------------------------------
# Main Verification Function
# ------------------------------------------------------------------

def verify_task(file_path):
    print(f"Verifying presentation: {file_path}")
    total_score = 0.0
    max_score = 1.0

    # Scoring weights (must sum to ≤ 1.0)
    slide_weight = 0.2    # Slide 91 exists
    table_weight = 0.4    # Correct-sized table exists
    align_weight = 0.4    # Column 2 right-aligned

    # ---------- Load presentation ----------
    if not os.path.isfile(file_path):
        print("✗ File does not exist")
        print("REWARD: 0.0")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error opening presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # ---------- Verify Slide 91 ----------
    slide_index = 90  # zero-based index for slide 91
    if len(prs.slides) > slide_index:
        print(f"✓ Slide 91 exists (total slides: {len(prs.slides)})")
        total_score += slide_weight
        slide = prs.slides[slide_index]
    else:
        print(f"✗ Slide 91 does not exist (total slides: {len(prs.slides)})")
        print(f"Total Score: {total_score} (out of {max_score})")
        print(f"REWARD: {total_score}")
        return total_score

    # ---------- Search for qualifying table ----------
    table_found = False
    alignment_ok = False

    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            print(f"  Found table with {len(table.rows)} rows and {len(table.columns)} columns")
            if len(table.rows) == 2 and len(table.columns) == 10:
                table_found = True
                if verify_table_alignment(table):
                    alignment_ok = True
                # break after checking the first 2×10 table (whether alignment passes or not)
                break

    # ---------- Scoring ----------
    if table_found:
        print("✓ Table with correct dimensions found on slide 91")
        total_score += table_weight
        if alignment_ok:
            print("✓ Text in column 2 is right-aligned for both rows")
            total_score += align_weight
        else:
            print("✗ Column 2 text is not properly right-aligned")
    else:
        print("✗ No 2×10 table found on slide 91")

    # ---------- Finalise score ----------
    final_score = min(total_score, max_score)
    print(f"Total Score: {final_score} (out of {max_score})")
    print(f"REWARD: {final_score}")
    return final_score

# ------------------------------------------------------------------
# Script Entry Point (executed in grading environment)
# ------------------------------------------------------------------
if __name__ == "__main__":
    FILE_PATH = (
        "/home/user/quick_question_on_slide_91_of_my_libreoffice_impress_deck_i_need_to_insert_a_table_thats_exactly_2_r_golden.pptx"
    )
    verify_task(FILE_PATH)

