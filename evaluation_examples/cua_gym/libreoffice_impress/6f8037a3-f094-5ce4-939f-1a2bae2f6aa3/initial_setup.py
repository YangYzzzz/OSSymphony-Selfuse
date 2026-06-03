"""
Initial Setup: Sales pitch deck with white backgrounds and black text (9 slides)
Task ID: osworld_impress_all_slides_background_011
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_all_slides_background_011'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_slide_background_white(slide):
    """Set slide background to white."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                color=RGBColor(0x00, 0x00, 0x00), alignment=PP_ALIGN.LEFT):
    """Add a textbox with black text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_text_lines(slide, left, top, width, height, lines, font_size=16,
                   color=RGBColor(0x00, 0x00, 0x00)):
    """Add multiple lines of text to a textbox."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    # Standard widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    BLACK = RGBColor(0x00, 0x00, 0x00)
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)

    # ---- Slide 1: Title / Cover ----
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide1)
    add_textbox(slide1, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
                "NexaTech Solutions", font_size=44, bold=True,
                color=BLACK, alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(1.5), Inches(3.7), Inches(10), Inches(1.0),
                "Transforming Business Through Intelligent Automation",
                font_size=22, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(1.5), Inches(5.0), Inches(10), Inches(0.6),
                "Q1 2025 Investor Presentation | Confidential",
                font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # ---- Slide 2: Executive Summary ----
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide2)
    add_textbox(slide2, Inches(0.7), Inches(0.4), Inches(11.5), Inches(0.9),
                "Executive Summary", font_size=32, bold=True, color=BLACK)
    add_text_lines(slide2, Inches(0.7), Inches(1.5), Inches(11.5), Inches(5.5), [
        "• NexaTech Solutions is a B2B SaaS company specializing in workflow automation",
        "• Founded in 2019, headquartered in Austin, TX with 148 full-time employees",
        "• ARR of $12.4M growing at 67% YoY — targeting $20M by end of 2025",
        "• Serving 312 enterprise clients across 14 industry verticals",
        "• Seeking Series B funding of $25M to accelerate product development and GTM",
        "• Strategic partnerships with Salesforce, Microsoft Azure, and SAP",
    ], font_size=18, color=BLACK)

    # ---- Slide 3: Problem Statement ----
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide3)
    add_textbox(slide3, Inches(0.7), Inches(0.4), Inches(11.5), Inches(0.9),
                "The Problem", font_size=32, bold=True, color=BLACK)
    add_textbox(slide3, Inches(0.7), Inches(1.4), Inches(11.5), Inches(0.7),
                "Enterprise teams lose 32% of productive hours to manual, repetitive workflows",
                font_size=20, color=DARK_GRAY)
    add_text_lines(slide3, Inches(0.7), Inches(2.3), Inches(11.5), Inches(4.5), [
        "  Manual data entry errors cost US businesses $3.1 trillion annually (IDC Research)",
        "  62% of IT leaders cite legacy process bottlenecks as top innovation barrier",
        "  Average enterprise uses 254 SaaS applications — 60% operate in silos",
        "  Cross-department handoffs average 4.7 days due to fragmented tooling",
        "  Compliance reporting consumes 18% of finance team bandwidth each quarter",
    ], font_size=17, color=BLACK)

    # ---- Slide 4: Our Solution ----
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide4)
    add_textbox(slide4, Inches(0.7), Inches(0.4), Inches(11.5), Inches(0.9),
                "Our Solution", font_size=32, bold=True, color=BLACK)
    add_textbox(slide4, Inches(0.7), Inches(1.4), Inches(11.5), Inches(0.7),
                "NexaFlow Platform — End-to-end intelligent process automation",
                font_size=20, color=DARK_GRAY)
    add_text_lines(slide4, Inches(0.7), Inches(2.3), Inches(11.5), Inches(4.5), [
        "  AI-Powered Workflow Builder: Drag-and-drop automation with ML-based suggestions",
        "  Universal Integration Hub: Native connectors for 800+ enterprise applications",
        "  Real-time Analytics Dashboard: Live KPI monitoring with anomaly detection",
        "  Compliance Automation Engine: Automated audit trails and regulatory reporting",
        "  Collaboration Layer: Role-based approvals, notifications, and task assignments",
        "  Low-code / No-code: Business users build workflows without engineering support",
    ], font_size=17, color=BLACK)

    # ---- Slide 5: Market Opportunity ----
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide5)
    add_textbox(slide5, Inches(0.7), Inches(0.4), Inches(11.5), Inches(0.9),
                "Market Opportunity", font_size=32, bold=True, color=BLACK)
    add_text_lines(slide5, Inches(0.7), Inches(1.5), Inches(5.5), Inches(5.0), [
        "Total Addressable Market",
        "$47.8B by 2027",
        "(BPA + RPA + iBPMS)",
        "",
        "Serviceable Addressable Market",
        "$8.2B mid-market enterprises",
        "",
        "Current Market Share",
        "0.15% — significant upside",
    ], font_size=18, color=BLACK)
    add_text_lines(slide5, Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.0), [
        "Key Growth Drivers:",
        "• Digital transformation mandates accelerating",
        "• Remote/hybrid work requiring process standardization",
        "• Regulatory complexity increasing automation demand",
        "• AI maturity enabling smarter automation ROI",
        "• Cloud-native adoption lowering deployment barriers",
    ], font_size=17, color=BLACK)

    # ---- Slide 6: Traction & Metrics ----
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide6)
    add_textbox(slide6, Inches(0.7), Inches(0.4), Inches(11.5), Inches(0.9),
                "Traction & Key Metrics", font_size=32, bold=True, color=BLACK)
    add_text_lines(slide6, Inches(0.7), Inches(1.5), Inches(5.5), Inches(5.0), [
        "Financial Performance:",
        "ARR: $12.4M (+67% YoY)",
        "MRR: $1.03M (Sep 2024)",
        "Gross Margin: 74%",
        "Net Revenue Retention: 118%",
        "CAC Payback Period: 14 months",
        "LTV:CAC Ratio: 5.8x",
    ], font_size=17, color=BLACK)
    add_text_lines(slide6, Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.0), [
        "Customer Success:",
        "312 enterprise clients",
        "Average ACV: $39,700",
        "Churn Rate: 3.2% annually",
        "NPS Score: 67",
        "Workflows Automated: 2.1M+",
        "Hours Saved Monthly: 840,000+",
    ], font_size=17, color=BLACK)

    # ---- Slide 7: Team ----
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide7)
    add_textbox(slide7, Inches(0.7), Inches(0.4), Inches(11.5), Inches(0.9),
                "Leadership Team", font_size=32, bold=True, color=BLACK)
    add_text_lines(slide7, Inches(0.7), Inches(1.5), Inches(5.5), Inches(5.5), [
        "Alexandra Rivera — CEO & Co-Founder",
        "Former VP Product at Workday; Stanford MBA",
        "Led $0 → $15M ARR at previous startup (acquired 2018)",
        "",
        "Jordan Kim — CTO & Co-Founder",
        "Ex-Principal Engineer at Stripe; MIT CS PhD",
        "Filed 7 patents in distributed systems",
        "",
        "Marcus Okonkwo — CFO",
        "Former CFO at PagerDuty; 3 IPO exits",
        "CPA, CFA, 22 years finance leadership",
    ], font_size=16, color=BLACK)
    add_text_lines(slide7, Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.5), [
        "Priya Nair — VP Sales",
        "Built $40M ARR sales team at Zendesk",
        "Expertise in enterprise mid-market GTM",
        "",
        "Daniel Ferreira — VP Engineering",
        "Previously at Databricks and Twilio",
        "Scaled engineering from 12 to 65 engineers",
        "",
        "Advisory Board:",
        "Susan Caldwell — Former CIO, General Electric",
        "Thomas Park — Partner, Sequoia Capital",
        "Elena Volkov — Board, Salesforce Ventures",
    ], font_size=16, color=BLACK)

    # ---- Slide 8: Roadmap ----
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide8)
    add_textbox(slide8, Inches(0.7), Inches(0.4), Inches(11.5), Inches(0.9),
                "Product Roadmap 2025", font_size=32, bold=True, color=BLACK)
    add_text_lines(slide8, Inches(0.7), Inches(1.5), Inches(3.7), Inches(5.5), [
        "Q1 2025 (In Progress)",
        "• NexaFlow 3.0 release",
        "• Advanced AI decision engine",
        "• Salesforce CRM deep integration",
        "• SOC 2 Type II certification",
    ], font_size=16, color=BLACK)
    add_text_lines(slide8, Inches(4.9), Inches(1.5), Inches(3.7), Inches(5.5), [
        "Q2-Q3 2025",
        "• Mobile app (iOS / Android)",
        "• ERP connector suite (SAP, Oracle)",
        "• Predictive workflow analytics",
        "• EU data residency compliance",
    ], font_size=16, color=BLACK)
    add_text_lines(slide8, Inches(9.1), Inches(1.5), Inches(3.7), Inches(5.5), [
        "Q4 2025",
        "• NexaFlow Marketplace launch",
        "• Third-party workflow templates",
        "• AI co-pilot beta (GPT-4 powered)",
        "• Enterprise SSO enhancements",
    ], font_size=16, color=BLACK)

    # ---- Slide 9: Call to Action / Ask ----
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide9)
    add_textbox(slide9, Inches(0.7), Inches(0.4), Inches(11.5), Inches(0.9),
                "The Ask & Use of Funds", font_size=32, bold=True, color=BLACK)
    add_textbox(slide9, Inches(0.7), Inches(1.4), Inches(11.5), Inches(0.7),
                "Raising $25M Series B at $120M pre-money valuation",
                font_size=20, color=DARK_GRAY)
    add_text_lines(slide9, Inches(0.7), Inches(2.4), Inches(5.5), Inches(4.5), [
        "Use of Funds:",
        "40% — Product R&D",
        "   AI/ML capabilities expansion",
        "   Platform scalability & security",
        "30% — Sales & Marketing",
        "   Enterprise GTM expansion",
        "   Partner channel development",
        "20% — Customer Success",
        "   Onboarding & support teams",
        "10% — G&A / Working Capital",
    ], font_size=16, color=BLACK)
    add_text_lines(slide9, Inches(7.0), Inches(2.4), Inches(5.5), Inches(4.5), [
        "Expected Milestones by EOY 2025:",
        "• ARR: $20M target (+61% growth)",
        "• Customer Count: 550+ enterprises",
        "• Team: Grow to 225 employees",
        "• Markets: Expand to EU & APAC",
        "• IPO readiness: Target 2027",
        "",
        "Contact:",
        "Alexandra Rivera, CEO",
        "a.rivera@nexatechsolutions.com",
        "+1 (512) 874-3300",
    ], font_size=16, color=BLACK)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
