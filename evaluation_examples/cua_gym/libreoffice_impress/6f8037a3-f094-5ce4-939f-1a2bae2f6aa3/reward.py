"""
Reward Script: Apply dark navy background (#1A237E) to all slides and white text on all slides
Task ID: osworld_impress_all_slides_background_011
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.5): All 9 slides have #1A237E background
  - Component 2 (0.3): All text runs across all slides are white (#FFFFFF)
  - Component 3 (0.2): Compound check — all 9 slides have BOTH correct background AND all-white text
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_all_slides_background_011'

EXPECTED_BG = '1A237E'   # dark navy
EXPECTED_TEXT = 'FFFFFF'  # white
EXPECTED_SLIDE_COUNT = 9


def get_slide_background_rgb(slide):
    """Return hex string of slide background color, or None if not a solid fill."""
    fill = slide.background.fill
    if fill.type == 1:  # solid fill
        try:
            return str(fill.fore_color.rgb).upper()
        except Exception:
            return None
    # type 5 = inherited — treat as no explicit bg set
    return None


def get_all_text_runs(slide):
    """Recursively collect all non-empty text runs from all shapes (including groups)."""
    runs = []

    def extract(shape):
        if hasattr(shape, 'text_frame') and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or '').strip():
                        runs.append(run)
        if hasattr(shape, 'shapes'):  # group shapes
            for sub in shape.shapes:
                extract(sub)

    for shape in slide.shapes:
        extract(shape)
    return runs


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    if num_slides != EXPECTED_SLIDE_COUNT:
        print(f"WARN: Expected {EXPECTED_SLIDE_COUNT} slides, found {num_slides}")

    # Component 1: All slides have #1A237E (dark navy) background (0.5 points)
    # This FAILS on initial (all FFFFFF backgrounds) → PASSES on golden (all 1A237E backgrounds)
    try:
        slides_with_correct_bg = 0
        bg_details = []
        for i, slide in enumerate(prs.slides):
            bg_color = get_slide_background_rgb(slide)
            if bg_color == EXPECTED_BG:
                slides_with_correct_bg += 1
            else:
                bg_details.append(f"Slide {i+1}: {bg_color}")

        if slides_with_correct_bg == EXPECTED_SLIDE_COUNT:
            print(f"PASS: Component 1 — All {EXPECTED_SLIDE_COUNT} slides have #1A237E background (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 1 — {slides_with_correct_bg}/{EXPECTED_SLIDE_COUNT} slides have correct background")
            if bg_details:
                print(f"  Wrong backgrounds: {bg_details[:5]}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: All text runs across all slides are white (#FFFFFF) (0.3 points)
    # This FAILS on initial (black/dark text) → PASSES on golden (all white text)
    try:
        total_runs = 0
        white_runs = 0
        non_white_examples = []

        for i, slide in enumerate(prs.slides):
            runs = get_all_text_runs(slide)
            for run in runs:
                total_runs += 1
                try:
                    color_type = run.font.color.type
                    if color_type is not None:
                        rgb = str(run.font.color.rgb).upper()
                        if rgb == EXPECTED_TEXT:
                            white_runs += 1
                        else:
                            non_white_examples.append(
                                f"Slide {i+1}: '{run.text[:15]}' → {rgb}"
                            )
                    else:
                        # No explicit color = inherited (not explicitly set to white)
                        non_white_examples.append(
                            f"Slide {i+1}: '{run.text[:15]}' → inherited (not explicitly white)"
                        )
                except Exception:
                    non_white_examples.append(
                        f"Slide {i+1}: '{run.text[:15]}' → color error"
                    )

        if total_runs > 0 and white_runs == total_runs:
            print(f"PASS: Component 2 — All {total_runs} text runs are white (#FFFFFF) (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 — {white_runs}/{total_runs} text runs are white (#FFFFFF)")
            if non_white_examples:
                print(f"  Non-white examples: {non_white_examples[:5]}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Compound check — every slide has BOTH correct background AND all-white text (0.2 points)
    # This ensures uniform application across ALL slides, not just some
    try:
        def run_is_white(run):
            """Return True iff the run has an explicit #FFFFFF color."""
            try:
                return run.font.color.type is not None and str(run.font.color.rgb).upper() == EXPECTED_TEXT
            except Exception:
                return False

        def slide_text_all_white(slide):
            """Return True iff every non-empty text run on the slide is explicitly white."""
            runs = get_all_text_runs(slide)
            if not runs:
                return True  # no text → passes
            return all(run_is_white(r) for r in runs)

        slides_fully_correct = sum(
            1
            for slide in prs.slides
            if get_slide_background_rgb(slide) == EXPECTED_BG and slide_text_all_white(slide)
        )

        if slides_fully_correct == EXPECTED_SLIDE_COUNT:
            print(f"PASS: Component 3 — All {EXPECTED_SLIDE_COUNT} slides fully correct (bg + text) (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 3 — Only {slides_fully_correct}/{EXPECTED_SLIDE_COUNT} slides have both correct bg and white text")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in a given env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
