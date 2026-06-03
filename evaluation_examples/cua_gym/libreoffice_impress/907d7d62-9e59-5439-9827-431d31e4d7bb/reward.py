"""
FINAL REWARD SCRIPT - SUCCESS
Task: In LibreOffice Impress, slide 248 needs a tweak: I want the title aligned flush left and an exact 0.3 cm space after the title paragraph. How do I set that up?
Generated: 2025-09-10 20:04:23
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT


def verify_impress_task(file_path: str) -> float:
    """
    Verify that on slide 248 of the provided presentation:
    1. A title shape exists
    2. The first paragraph of the title is left-aligned
    3. The paragraph has exactly (≈) 0.3 cm space after

    Returns a progressive score between 0.0 and 1.0.
    """

    print("Starting verification for task: slide 248 title formatting")
    score = 0.0  # progressive score
    max_score = 1.0

    # --- Load presentation ---
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0

    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded presentation with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        return 0.0

    # --- Requirement 1: Slide 248 exists & has title ---
    if len(prs.slides) < 248:
        print(f"✗ Presentation has only {len(prs.slides)} slides – slide 248 missing")
        return 0.0

    slide = prs.slides[247]  # zero-indexed

    # Locate the title shape (prefer real title placeholder)
    title_shape = None
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # Prefer placeholder of type TITLE (value 1)
        if hasattr(shape, "placeholder_format"):
            try:
                if shape.placeholder_format.type == 1:  # TITLE
                    title_shape = shape
                    break
            except Exception:
                pass
        # Fallback: name contains "title"
        if title_shape is None and "title" in shape.name.lower():
            title_shape = shape

    if title_shape is None:
        print("✗ No title shape found on slide 248")
        return 0.0

    print("✓ Title shape found (0.2 points)")
    score += 0.2

    # --- Inspect the first paragraph of the title ---
    tf = title_shape.text_frame
    if not tf.paragraphs:
        print("✗ Title text frame has no paragraphs")
        return score

    para = tf.paragraphs[0]

    # Requirement 2: Left alignment
    if para.alignment == PP_PARAGRAPH_ALIGNMENT.LEFT:
        print("✓ Title paragraph aligned LEFT (0.4 points)")
        score += 0.4
    else:
        print(f"✗ Alignment is {para.alignment}, expected LEFT")

    # Requirement 3: Space-after ≈ 0.3 cm (0.4 points)
    # 1 cm = 360 000 EMUs in pptx, so 0.3 cm ≈ 108 000 EMUs
    expected_emus = int(360000 * 0.3)  # 108000
    actual_emus = para.space_after  # may be None

    if actual_emus is None:
        print("✗ space_after not set on paragraph")
    else:
        tolerance = 2000  # allow small rounding tolerance (~0.006 cm)
        if abs(actual_emus - expected_emus) <= tolerance:
            print(f"✓ space_after ≈ 0.3 cm (actual {actual_emus}) (0.4 points)")
            score += 0.4
        else:
            print(f"✗ space_after {actual_emus} EMUs ≠ 0.3 cm (expected {expected_emus})")

    final_score = min(score, max_score)
    print(f"Final Score: {final_score}")
    return final_score


if __name__ == "__main__":
    # Path to the presentation inside the VM environment
    FILE_PATH = "/home/user/in_libreoffice_impress_slide_248_needs_a_tweak_i_want_the_title_aligned_flush_left_and_an_exact_03_c_golden.pptx"

    reward = verify_impress_task(FILE_PATH)
    print(f"REWARD: {reward}")

