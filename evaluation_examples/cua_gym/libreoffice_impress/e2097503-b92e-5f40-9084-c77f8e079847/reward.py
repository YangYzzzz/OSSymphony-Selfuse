"""
Reward Script: Add horizontal line divider on slide 4 with 'Before' above and 'After' below
Task ID: impress_tm_085
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Horizontal line/divider shape on slide 4 near vertical center
  Component 2 (0.3): 'Before' text box positioned above the divider
  Component 3 (0.3): 'After' text box positioned below the divider
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_085'


def find_divider_shape(slide, slide_height):
    """
    Find a horizontal line/divider shape on the slide.
    A divider is a shape that is much wider than tall (aspect ratio > 10:1)
    and positioned near the vertical center of the slide.
    Could be an AUTO_SHAPE (thin rectangle), a LINE, or FREEFORM.
    """
    slide_mid_y = slide_height / 2
    candidates = []
    for shape in slide.shapes:
        # Skip shapes with substantial text content (titles, text boxes with content)
        if hasattr(shape, 'text') and shape.text.strip():
            continue
        w = shape.width
        h = shape.height
        # A divider should be wide and very thin (aspect ratio >= 10)
        if w > 0 and h >= 0 and h < w / 5:
            # Check the shape center is reasonably near vertical middle
            shape_center_y = shape.top + h / 2
            # Allow within 30% of slide height from center
            if abs(shape_center_y - slide_mid_y) < slide_height * 0.30:
                candidates.append(shape)
    return candidates


def find_text_shape(slide, target_text):
    """
    Find a text shape containing the target text (case-insensitive).
    Returns the shape if found, None otherwise.
    """
    target_lower = target_text.lower().strip()
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip().lower() == target_lower:
            return shape
        # Also check inside group shapes
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                if hasattr(sub, 'text') and sub.text.strip().lower() == target_lower:
                    return sub
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)
    slide_height = prs.slide_height
    slide_width = prs.slide_width

    # Component 1: Horizontal line/divider shape on slide 4 near vertical center (0.4 points)
    divider_shape = None
    try:
        dividers = find_divider_shape(slide, slide_height)
        if dividers:
            divider_shape = dividers[0]
            # Verify it spans a meaningful portion of the slide width (at least 50%)
            if divider_shape.width >= slide_width * 0.5:
                print(f"PASS: Component 1 — Horizontal divider found: "
                      f"pos=({divider_shape.left},{divider_shape.top}), "
                      f"size=({divider_shape.width},{divider_shape.height}), "
                      f"name='{divider_shape.name}' (0.4 pts)")
                total_score += 0.4
            else:
                print(f"FAIL: Component 1 — Shape found but too narrow: "
                      f"width={divider_shape.width}, need >= {slide_width * 0.5}")
        else:
            print(f"FAIL: Component 1 — No horizontal divider shape found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: 'Before' text box positioned above the divider (0.3 points)
    try:
        before_shape = find_text_shape(slide, 'Before')
        if before_shape is not None:
            if divider_shape is not None:
                divider_top = divider_shape.top
                before_bottom = before_shape.top + before_shape.height
                # 'Before' text box should be above the divider line
                if before_shape.top < divider_top:
                    print(f"PASS: Component 2 — 'Before' text found above divider: "
                          f"before_top={before_shape.top}, divider_top={divider_top} (0.3 pts)")
                    total_score += 0.3
                else:
                    print(f"FAIL: Component 2 — 'Before' text found but NOT above divider: "
                          f"before_top={before_shape.top}, divider_top={divider_top}")
            else:
                # No divider found, but 'Before' text exists in upper half of slide
                if before_shape.top < slide_height / 2:
                    print(f"PASS: Component 2 — 'Before' text found in upper half (no divider ref) (0.3 pts)")
                    total_score += 0.3
                else:
                    print(f"FAIL: Component 2 — 'Before' text found but in lower half")
        else:
            print(f"FAIL: Component 2 — No 'Before' text found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: 'After' text box positioned below the divider (0.3 points)
    try:
        after_shape = find_text_shape(slide, 'After')
        if after_shape is not None:
            if divider_shape is not None:
                divider_bottom = divider_shape.top + divider_shape.height
                # 'After' text box should be below the divider line
                if after_shape.top > divider_shape.top:
                    print(f"PASS: Component 3 — 'After' text found below divider: "
                          f"after_top={after_shape.top}, divider_bottom={divider_bottom} (0.3 pts)")
                    total_score += 0.3
                else:
                    print(f"FAIL: Component 3 — 'After' text found but NOT below divider: "
                          f"after_top={after_shape.top}, divider_top={divider_shape.top}")
            else:
                # No divider found, but 'After' text exists in lower half
                if after_shape.top >= slide_height / 2:
                    print(f"PASS: Component 3 — 'After' text found in lower half (no divider ref) (0.3 pts)")
                    total_score += 0.3
                else:
                    print(f"FAIL: Component 3 — 'After' text found but in upper half")
        else:
            print(f"FAIL: Component 3 — No 'After' text found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
