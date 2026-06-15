"""
Initial Setup: Create a 6-slide presentation with slide 4 titled 'Transformation' (title only).
Task ID: impress_tm_085
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_085'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    slide_width = prs.slide_width   # default 10 inches
    slide_height = prs.slide_height  # default 7.5 inches

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Project Comparison Report"
    slide1.placeholders[1].text = "Q1 2025 Strategic Review"

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "This presentation examines key metrics across departments"
    p2 = body2.add_paragraph()
    p2.text = "Revenue growth tracked against industry benchmarks"
    p2.level = 1
    p3 = body2.add_paragraph()
    p3.text = "Customer satisfaction scores by region"
    p3.level = 1
    p4 = body2.add_paragraph()
    p4.text = "Operational efficiency improvements since Q3 2024"
    p4.level = 1

    # --- Slide 3: Methodology ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Methodology"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Data collected from 12 regional offices over 6 months"
    p5 = body3.add_paragraph()
    p5.text = "Survey responses from 2,400 employees"
    p5.level = 1
    p6 = body3.add_paragraph()
    p6.text = "Financial data validated against audited reports"
    p6.level = 1
    p7 = body3.add_paragraph()
    p7.text = "Statistical analysis using paired t-tests and ANOVA"
    p7.level = 1

    # --- Slide 4: Transformation (Title Only - NO line, NO Before/After) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add just a title text box at the top
    title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Transformation"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True

    # --- Slide 5: Results ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Results"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Revenue increased by 18% compared to the prior year"
    p8 = body5.add_paragraph()
    p8.text = "Customer NPS improved from 42 to 67 across all regions"
    p8.level = 1
    p9 = body5.add_paragraph()
    p9.text = "Employee engagement scores rose by 15 points"
    p9.level = 1

    # --- Slide 6: Conclusions ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Conclusions"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Strategic initiatives delivered measurable impact"
    p10 = body6.add_paragraph()
    p10.text = "Recommend expanding regional pilot programs in Q2 2025"
    p10.level = 1
    p11 = body6.add_paragraph()
    p11.text = "Continue investment in employee development and training"
    p11.level = 1

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
