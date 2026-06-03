"""
Reward Script: Add superscript equation text box to slide 6
Task ID: impress_stu_043
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.25): New text box exists on slide 6 with equation text 'E = mc2'
  - Component 2 (0.25): Superscript '2' (baseline > 0 on the '2' run)
  - Component 3 (0.25): Font is Times New Roman, size 36pt on all runs
  - Component 4 (0.25): Paragraph alignment is centered
"""

import os
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_043'


def persist_app_state(domain):
    """Save any unsaved GUI state before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def find_equation_textbox(slide):
    """Find a text box on slide 6 that contains the equation text 'E = mc2'."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shape.has_text_frame:
            full_text = ""
            for para in shape.text_frame.paragraphs:
                full_text += para.text
            # Normalize whitespace for comparison
            normalized = full_text.replace(" ", "").lower()
            if "e=mc2" in normalized:
                return shape
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide6 = prs.slides[5]  # 0-indexed

    # Component 1: New text box exists on slide 6 with equation text containing 'E = mc2' (0.25 pts)
    try:
        eq_shape = find_equation_textbox(slide6)
        if eq_shape is not None:
            print(f"PASS: Component 1 — Equation text box found with text '{eq_shape.text_frame.paragraphs[0].text}' (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — No text box with 'E = mc2' equation found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    if eq_shape is None:
        # No equation text box found; remaining checks cannot proceed
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Get runs from the equation text box
    eq_para = eq_shape.text_frame.paragraphs[0]
    eq_runs = [r for r in eq_para.runs if (r.text or "").strip()]

    # Component 2: Superscript '2' — the '2' run has baseline > 0 (0.25 pts)
    try:
        found_superscript_2 = False
        for run in eq_runs:
            if "2" in run.text.strip():
                baseline = run.font._element.attrib.get('baseline', None)
                if baseline is not None and int(baseline) > 0:
                    found_superscript_2 = True
                    print(f"PASS: Component 2 — Superscript '2' found (baseline={baseline}) (0.25 pts)")
                    total_score += 0.25
                else:
                    print(f"FAIL: Component 2 — '2' run found but baseline={baseline}, expected positive value for superscript")
                break
        if not found_superscript_2 and total_score < 0.5:
            print(f"FAIL: Component 2 — No run containing '2' with superscript baseline found")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Font is Times New Roman, size 36pt on all runs (0.25 pts)
    try:
        all_correct = True
        expected_size = Pt(36)  # 457200 EMU
        for run in eq_runs:
            font_name = run.font.name
            font_size = run.font.size
            if font_name != "Times New Roman":
                print(f"FAIL: Component 3 — Run '{run.text}' has font '{font_name}', expected 'Times New Roman'")
                all_correct = False
                break
            if font_size != expected_size:
                print(f"FAIL: Component 3 — Run '{run.text}' has size {font_size}, expected {expected_size} (36pt)")
                all_correct = False
                break
        if all_correct and len(eq_runs) > 0:
            print(f"PASS: Component 3 — All runs have Times New Roman 36pt (0.25 pts)")
            total_score += 0.25
        elif len(eq_runs) == 0:
            print(f"FAIL: Component 3 — No non-empty runs found in equation text box")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Paragraph alignment is centered (0.25 pts)
    try:
        alignment = eq_para.alignment
        if alignment == PP_ALIGN.CENTER:
            print(f"PASS: Component 4 — Paragraph alignment is CENTER (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 4 — Alignment is {alignment}, expected CENTER (2)")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
