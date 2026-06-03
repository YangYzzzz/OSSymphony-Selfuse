"""
Initial Setup: Physics Relativity presentation with 8 slides.
Task ID: impress_stu_043
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_043'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Physics of Relativity"
    slide1.placeholders[1].text = "An Introduction to Einstein's Theories"

    # --- Slide 2: Historical Context ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Historical Context"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "In the early 20th century, physics faced several puzzles"
    p = tf2.add_paragraph()
    p.text = "The Michelson-Morley experiment (1887) found no evidence of luminiferous aether"
    p.level = 1
    p = tf2.add_paragraph()
    p.text = "Maxwell's equations predicted a constant speed of light"
    p.level = 1
    p = tf2.add_paragraph()
    p.text = "Newtonian mechanics could not explain Mercury's orbital precession"
    p.level = 1

    # --- Slide 3: Special Relativity ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Special Relativity (1905)"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Two fundamental postulates:"
    p = tf3.add_paragraph()
    p.text = "1. The laws of physics are the same in all inertial reference frames"
    p.level = 1
    p = tf3.add_paragraph()
    p.text = "2. The speed of light in vacuum is constant for all observers"
    p.level = 1
    p = tf3.add_paragraph()
    p.text = "Consequences: time dilation, length contraction, relativity of simultaneity"
    p.level = 0

    # --- Slide 4: Time Dilation ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Time Dilation"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Moving clocks run slower relative to a stationary observer"
    p = tf4.add_paragraph()
    p.text = "GPS satellites must account for relativistic time corrections"
    p.level = 1
    p = tf4.add_paragraph()
    p.text = "Muons created in the upper atmosphere survive to reach Earth's surface"
    p.level = 1
    p = tf4.add_paragraph()
    p.text = "The twin paradox illustrates asymmetric aging between travelers"
    p.level = 1

    # --- Slide 5: General Relativity ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "General Relativity (1915)"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Gravity is the curvature of spacetime caused by mass and energy"
    p = tf5.add_paragraph()
    p.text = "Predicted gravitational lensing, confirmed during the 1919 solar eclipse"
    p.level = 1
    p = tf5.add_paragraph()
    p.text = "Explained Mercury's perihelion precession with remarkable accuracy"
    p.level = 1
    p = tf5.add_paragraph()
    p.text = "Predicted the existence of gravitational waves, detected by LIGO in 2015"
    p.level = 1

    # --- Slide 6: Mass-Energy Equivalence (task target) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Mass-Energy Equivalence"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "One of the most famous consequences of special relativity"
    p = tf6.add_paragraph()
    p.text = "Mass and energy are interchangeable quantities"
    p.level = 1
    p = tf6.add_paragraph()
    p.text = "A small amount of mass contains an enormous amount of energy"
    p.level = 1
    p = tf6.add_paragraph()
    p.text = "This principle underlies nuclear fission and fusion reactions"
    p.level = 1
    p = tf6.add_paragraph()
    p.text = "The sun converts approximately 4 million tons of mass to energy each second"
    p.level = 1
    # NO extra text box with E=mc^2 here — that is the task

    # --- Slide 7: Applications ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Real-World Applications"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "Relativity impacts many technologies we use daily"
    p = tf7.add_paragraph()
    p.text = "GPS navigation requires relativistic corrections for accurate positioning"
    p.level = 1
    p = tf7.add_paragraph()
    p.text = "Particle accelerators like CERN rely on relativistic mechanics"
    p.level = 1
    p = tf7.add_paragraph()
    p.text = "Nuclear power plants harness mass-energy equivalence"
    p.level = 1
    p = tf7.add_paragraph()
    p.text = "PET scanners in medicine use positron-electron annihilation"
    p.level = 1

    # --- Slide 8: Summary & Further Reading ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Summary & Further Reading"
    tf8 = slide8.placeholders[1].text_frame
    tf8.text = "Key takeaways from this presentation"
    p = tf8.add_paragraph()
    p.text = "Special relativity: constant speed of light, time dilation, length contraction"
    p.level = 1
    p = tf8.add_paragraph()
    p.text = "General relativity: gravity as spacetime curvature"
    p.level = 1
    p = tf8.add_paragraph()
    p.text = "Recommended: 'Relativity: The Special and the General Theory' by A. Einstein"
    p.level = 1
    p = tf8.add_paragraph()
    p.text = "Online resource: Stanford Encyclopedia of Philosophy — Spacetime"
    p.level = 1

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
