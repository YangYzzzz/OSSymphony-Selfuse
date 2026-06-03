"""
FINAL REWARD SCRIPT - SUCCESS
Task: Add caption 'Table 2: Metrics' below Table 2.
Generated: 2025-10-17 17:34:54
Status: success
Model: azure-o3
Total Steps: 3
"""

# Reward script for verifying that the caption "Table 2: Metrics" was added
# below Table 2 in the supplied presentation.
#
# Scoring (progressive):
#   0.6 points – Correct caption text is present anywhere in the file
#   0.4 points – Caption is positioned BELOW a table on the same slide
# The script prints detailed diagnostics and finally prints
#   REWARD: <score>
# where <score> is a float between 0.0 and 1.0.

from pptx import Presentation
import os

CAPTION_TEXT = "table 2: metrics"  # lower-case for case-insensitive match
EMU_TOLERANCE = 70000  # ~1 mm – tolerance for minor placement differences


def verify_task(file_path: str) -> float:
    """Verify that the caption 'Table 2: Metrics' exists and is placed below Table 2."""

    max_score = 1.0
    score = 0.0

    # ---------- 1.  Load presentation (no points for loading itself) ----------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded presentation with {len(prs.slides)} slides")
    except Exception as exc:
        print(f"✗ Error opening presentation: {exc}")
        print("REWARD: 0.0")
        return 0.0

    # ---------- 2.  Search for caption text & evaluate placement ----------
    caption_found = False        # did we find the exact caption text?
    placement_correct = False    # is it actually below a table?

    for slide_idx, slide in enumerate(prs.slides, start=1):
        # Collect tables and candidate caption shapes on this slide
        tables = [sh for sh in slide.shapes if sh.has_table]
        captions = [sh for sh in slide.shapes
                    if sh.has_text_frame and sh.text and CAPTION_TEXT in sh.text.strip().lower()]

        if captions:
            caption_found = True  # we have at least one occurrence of the caption text

        # If no tables or captions on this slide, skip placement evaluation
        if not tables or not captions:
            continue

        # Check each caption against each table for below-table placement
        for table in tables:
            table_bottom = table.top + table.height  # bottom edge of the table in EMUs
            for cap in captions:
                # Caption considered correct if its TOP is below (>=) table bottom (with tolerance)
                if cap.top + EMU_TOLERANCE >= table_bottom:
                    placement_correct = True
                    print(f"✓ Caption below table on slide {slide_idx}")
                    break
            if placement_correct:
                break  # no need to check further tables
        if placement_correct:
            break      # overall placement already confirmed

    # ---------- 3.  Progressive scoring ----------
    if caption_found:
        score += 0.6
        print("✓ Caption text found (0.6 points)")
    else:
        print("✗ Caption text 'Table 2: Metrics' not found (0 points)")

    if placement_correct:
        score += 0.4
        print("✓ Caption placement correct (0.4 points)")
    else:
        print("✗ Caption is not correctly placed below any table (0 points)")

    # Clamp final score to [0.0, 1.0]
    final_score = min(score, max_score)
    print(f"REWARD: {final_score}")
    return final_score


# --------------------  Run verification  --------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/add_caption_table_2_metrics_below_table_2.pptx"
    verify_task(FILE_PATH)

