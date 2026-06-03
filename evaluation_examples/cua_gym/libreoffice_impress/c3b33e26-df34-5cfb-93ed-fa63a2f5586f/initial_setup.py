"""
Initial Setup: Add entrance animations to slide 1 shapes
Task ID: impress_rp_004
Domain: libreoffice_impress

Creates a presentation with 3 shapes on slide 1 (title, subtitle, image).
No animations are applied.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_004'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
IMG_PATH = f'{WORKDIR}/product_image.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_product_image():
    """Create a realistic product image placeholder."""
    img = Image.new('RGB', (800, 600), color=(240, 240, 245))
    draw = ImageDraw.Draw(img)

    # Draw a stylized product box
    # Outer box
    draw.rounded_rectangle([150, 80, 650, 520], radius=20, fill=(45, 62, 80), outline=(52, 73, 94), width=3)
    # Inner screen area
    draw.rounded_rectangle([180, 110, 620, 420], radius=10, fill=(236, 240, 241))
    # Product detail lines
    draw.rectangle([200, 140, 400, 160], fill=(52, 152, 219))
    draw.rectangle([200, 180, 500, 195], fill=(189, 195, 199))
    draw.rectangle([200, 210, 480, 225], fill=(189, 195, 199))
    draw.rectangle([200, 240, 520, 255], fill=(189, 195, 199))
    # A colored accent bar
    draw.rectangle([200, 290, 600, 310], fill=(46, 204, 113))
    # Bottom label area
    draw.rectangle([250, 450, 550, 490], fill=(52, 73, 94))

    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 28)
        draw.text((280, 452), "ProMax 5000", fill=(255, 255, 255), font=font)
    except (OSError, IOError):
        draw.text((310, 460), "ProMax 5000", fill=(255, 255, 255))

    img.save(IMG_PATH)
    print(f'Product image created: {IMG_PATH}')


def create_initial():
    create_product_image()

    prs = Presentation()
    # Use blank layout for full control
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    # --- Title text box at top ---
    title_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(0.5), Inches(8.0), Inches(1.2)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Introducing ProMax 5000"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    # Name the shape for identification
    title_box.name = "Title"

    # --- Subtitle text box below title ---
    subtitle_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(1.8), Inches(7.0), Inches(0.8)
    )
    tf2 = subtitle_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "The Future of Productivity"
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.name = "Arial"
    run2.font.size = Pt(24)
    run2.font.italic = True
    run2.font.color.rgb = RGBColor(0x4A, 0x4A, 0x5A)
    subtitle_box.name = "Subtitle"

    # --- Product image centered in lower half ---
    pic = slide.shapes.add_picture(
        IMG_PATH,
        Inches(2.0), Inches(3.0),
        Inches(6.0), Inches(4.0)
    )
    pic.name = "ProductImage"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
