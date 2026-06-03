"""
FINAL REWARD SCRIPT - SUCCESS
Task: Replace all double spaces in the document with a single space.
Generated: 2025-10-17 07:47:44
Status: success
Model: azure-o3
Total Steps: 1
"""

import os
import re
import zipfile
from pptx import Presentation
from lxml import etree

def _collect_slide_texts(prs):
    """Return a list of (location, identifier, text) tuples for slide texts."""
    texts = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                txt = shape.text or ""
                texts.append(("slide", f"{slide_idx}", txt))
    return texts


def _collect_notes_texts(pptx_path):
    """Return a list of (location, identifier, text) tuples for notes texts via raw XML."""
    texts = []
    try:
        with zipfile.ZipFile(pptx_path, "r") as zf:
            notes_files = [f for f in zf.namelist() if f.startswith("ppt/notesSlides/") and f.endswith(".xml")]
            for nf in notes_files:
                xml_bytes = zf.read(nf)
                root = etree.fromstring(xml_bytes)
                # grab every <a:t> element (namespace‐aware)
                t_elems = root.xpath("//a:t", namespaces={"a": "http://schemas.openxmlformats.org/drawingml/2006/main"})
                combined = " ".join(t.text for t in t_elems if t is not None and t.text)
                if combined:
                    texts.append(("note", nf, combined))
    except Exception as exc:
        print(f"✗ Error extracting notes text: {exc}")
    return texts


def verify_task(pptx_path):
    """Verify that ALL texts in the presentation contain no double (or more) spaces.

    Scoring:
      • 1.0  – every text string (slides + notes) is free of consecutive spaces.
      • <1.0 – proportional to the share of text strings without double spaces.
      • 0.0  – file missing / cannot be opened.
    """
    if not os.path.exists(pptx_path):
        print(f"✗ File not found: {pptx_path}")
        return 0.0

    # Try loading the PPTX
    try:
        prs = Presentation(pptx_path)
    except Exception as exc:
        print(f"✗ Unable to load PPTX: {exc}")
        return 0.0

    # Gather all text entries from slides and notes
    texts = _collect_slide_texts(prs) + _collect_notes_texts(pptx_path)

    # If presentation has no text at all, the requirement is trivially satisfied
    if not texts:
        print("✓ No text found – no double spaces present by definition")
        print("REWARD: 1.0")
        return 1.0

    double_space_regex = re.compile(r" {2,}")

    total_items = len(texts)
    bad_items = 0

    for loc, ident, txt in texts:
        if double_space_regex.search(txt):
            bad_items += 1
            sample = txt.replace("\n", " ")
            print(f"✗ Double spaces found in {loc} {ident}: '{sample[:60]}'")

    good_items = total_items - bad_items
    score = good_items / total_items if total_items else 1.0

    print(f"Total text elements checked: {total_items}")
    print(f"Elements WITHOUT double spaces: {good_items}")
    print(f"Verification score: {score}")

    print(f"REWARD: {score}")
    return score


if __name__ == "__main__":
    # Path provided in task context
    file_path = "/home/user/replace_all_double_spaces_in_the_document_with_a_single_space.pptx"
    verify_task(file_path)

