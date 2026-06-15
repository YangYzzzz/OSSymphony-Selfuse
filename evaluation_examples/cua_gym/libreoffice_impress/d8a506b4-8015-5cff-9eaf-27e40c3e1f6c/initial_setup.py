"""
Initial Setup: Create Academic Review presentation with 7 slides.
Task ID: impress_stu_037
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_037'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_body(slide, title_text, body_lines):
    """Helper to set title and body content on a slide."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    # Find content placeholder (index 1 typically)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            for i, line in enumerate(body_lines):
                if i == 0:
                    tf.paragraphs[0].text = line
                else:
                    p = tf.add_paragraph()
                    p.text = line
            break


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Academic Review 2025"
    slide1.placeholders[1].text = "Westfield Academy — Annual Performance Report"

    # --- Slide 2: Introduction ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide2, "Introduction", [
        "This presentation provides a comprehensive review of student academic performance",
        "Covering the 2024-2025 academic year across all STEM departments",
        "Data collected from midterm and final examinations",
        "Analysis includes grade distributions and subject-specific trends",
    ])

    # --- Slide 3: Course Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide3, "Course Overview", [
        "Mathematics: Advanced Calculus, Linear Algebra, Statistics",
        "Physics: Classical Mechanics, Electromagnetism, Thermodynamics",
        "Chemistry: Organic Chemistry, Analytical Chemistry, Biochemistry",
        "Biology: Molecular Biology, Genetics, Ecology",
    ])

    # --- Slide 4: Student Demographics ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide4, "Student Demographics", [
        "Total enrolled students: 347",
        "Senior year: 89 students (25.6%)",
        "Junior year: 112 students (32.3%)",
        "Sophomore year: 146 students (42.1%)",
        "Gender distribution: 52% female, 48% male",
    ])

    # --- Slide 5: Performance Overview (NO chart - task will add it) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    # Add title text box manually
    from pptx.util import Emu
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Performance Overview"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 6: Recommendations ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide6, "Recommendations", [
        "Increase tutoring support for Chemistry (lowest midterm average: 71)",
        "Continue Physics mentoring program (consistent performance)",
        "Expand Biology lab sessions given strong final results (92 avg)",
        "Introduce peer study groups for Mathematics",
    ])

    # --- Slide 7: Summary & Next Steps ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide7, "Summary & Next Steps", [
        "Overall improvement trend observed from midterm to final scores",
        "Chemistry showed the greatest improvement (+9 points average)",
        "Next review scheduled for September 2025",
        "Action items to be distributed to department heads by May 15",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
