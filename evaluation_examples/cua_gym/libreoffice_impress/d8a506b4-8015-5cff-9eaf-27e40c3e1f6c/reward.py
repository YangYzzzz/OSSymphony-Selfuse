"""
Reward Script: Add horizontal grouped bar chart on slide 5 comparing midterm vs final scores
Task ID: impress_stu_037
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): Chart exists on slide 5 with BAR_CLUSTERED type
  Component 2 (0.15): Chart title is 'Academic Performance Comparison'
  Component 3 (0.20): Correct categories: Math, Physics, Chemistry, Biology
  Component 4 (0.20): Correct Midterm series data: 78, 82, 71, 88
  Component 5 (0.15): Correct Final series data: 85, 79, 80, 92
  Component 6 (0.10): Two series with names Midterm and Final
"""

import os
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_037'


def persist_app_state():
    """Try to save any unsaved LibreOffice state via Ctrl+S."""
    try:
        os.environ["DISPLAY"] = ":0"
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        import time
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"PRECONDITION FAIL: Need at least 5 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide5 = prs.slides[4]  # 0-indexed

    # Find chart shape on slide 5
    chart_shape = None
    for shape in slide5.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Chart exists on slide 5 with BAR_CLUSTERED type (0.20 points)
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            chart_type_val = chart.chart_type
            # BAR_CLUSTERED = 57 in python-pptx (horizontal bars, grouped)
            if chart_type_val == 57:
                print(f"PASS: Component 1 -- BAR_CLUSTERED chart found on slide 5 (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 1 -- Chart found but type is {chart_type_val}, expected BAR_CLUSTERED (57)")
        else:
            print("FAIL: Component 1 -- No chart found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # If no chart at all, remaining checks will fail, return early
    if chart_shape is None:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    chart = chart_shape.chart

    # Component 2: Chart title is 'Academic Performance Comparison' (0.15 points)
    try:
        if chart.has_title:
            title_text = chart.chart_title.text_frame.text.strip()
            if title_text == 'Academic Performance Comparison':
                print(f"PASS: Component 2 -- Chart title matches exactly (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 2 -- Chart title is '{title_text}', expected 'Academic Performance Comparison'")
        else:
            print("FAIL: Component 2 -- Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Correct categories: Math, Physics, Chemistry, Biology (0.20 points)
    try:
        plot = chart.plots[0]
        categories = list(plot.categories)
        expected_categories = ['Math', 'Physics', 'Chemistry', 'Biology']
        if categories == expected_categories:
            print(f"PASS: Component 3 -- Categories match exactly: {categories} (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 3 -- Categories are {categories}, expected {expected_categories}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Get series names from XML for Components 4, 5, 6
    series_names = []
    try:
        ns = {'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'}
        chart_xml = chart._chartSpace.xml
        root = ET.fromstring(chart_xml)
        for ser in root.findall('.//c:ser', ns):
            tx = ser.find('.//c:tx//c:v', ns)
            if tx is not None:
                series_names.append(tx.text)
            else:
                series_names.append(None)
    except Exception as e:
        print(f"WARN: Could not parse series names from XML: {e}")

    # Component 4: Correct Midterm series data (0.20 points)
    try:
        plot = chart.plots[0]
        expected_midterm = [78.0, 82.0, 71.0, 88.0]
        all_series_vals = [list(s.values) for s in plot.series]
        if expected_midterm in all_series_vals:
            print(f"PASS: Component 4 -- Midterm data matches: {expected_midterm} (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 4 -- No series with midterm data {expected_midterm}. Found: {all_series_vals}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Correct Final series data (0.15 points)
    try:
        plot = chart.plots[0]
        expected_final = [85.0, 79.0, 80.0, 92.0]
        all_series_vals = [list(s.values) for s in plot.series]
        if expected_final in all_series_vals:
            print(f"PASS: Component 5 -- Final data matches: {expected_final} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 5 -- No series with final data {expected_final}. Found: {all_series_vals}")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # Component 6: Two series with names 'Midterm' and 'Final' (0.10 points)
    try:
        plot = chart.plots[0]
        num_series = len(plot.series)
        if num_series == 2 and set(series_names) == {'Midterm', 'Final'}:
            print(f"PASS: Component 6 -- 2 series named {series_names} (0.10 pts)")
            total_score += 0.10
        elif num_series == 2:
            print(f"FAIL: Component 6 -- 2 series but names are {series_names}, expected {{'Midterm', 'Final'}}")
        else:
            print(f"FAIL: Component 6 -- {num_series} series, expected 2")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
