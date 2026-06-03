"""
FINAL REWARD SCRIPT - SUCCESS
Task: Set the header row of Table 1 to bold and centered text.
Generated: 2025-10-17 07:23:24
Status: success
Model: azure-o3
Total Steps: 1
"""

import os
from pptx import Presentation
from pptx.enum.text import PP_ALIGN


def verify_header_bold_center(file_path: str) -> float:
    """Verify that the header row (row 0) of the first table in the
    presentation is entirely bold and center-aligned.

    Returns a progressive score between 0.0 and 1.0:
      • +0.5 if every header cell is bold
      • +0.5 if every header cell is horizontally centered
    """
    max_score = 1.0
    score = 0.0

    if not os.path.exists(file_path):
        print("✗ File not found:", file_path)
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Unable to load PPTX file: {e}")
        return 0.0

    # Locate the first table ("Table 1")
    table = None
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if shape.has_table:
                table = shape.table
                print(f"✓ Found table on slide {slide_idx}, shape {shape_idx}")
                break
        if table:
            break

    if table is None:
        print("✗ No table detected in presentation")
        return 0.0

    if len(table.rows) == 0:
        print("✗ Table contains no rows")
        return 0.0

    header_row = table.rows[0]
    num_cells = len(header_row.cells)
    print(f"Header row contains {num_cells} cells")

    bold_ok = True
    center_ok = True

    # Examine each cell in header row
    for cell_idx, cell in enumerate(header_row.cells, start=1):
        tf = cell.text_frame

        # Empty cell text means failure for both checks
        if not tf.text.strip():
            print(f"  Cell {cell_idx}: empty text – bold ✗, center ✗")
            bold_ok = center_ok = False
            continue

        # Check bold formatting across all runs in all paragraphs
        cell_bold = all(
            run.font.bold is True
            for para in tf.paragraphs
            for run in para.runs
        )

        # Check horizontal centering for all paragraphs
        cell_center = all(
            para.alignment == PP_ALIGN.CENTER for para in tf.paragraphs
        )

        print(
            f"  Cell {cell_idx}: bold={'✓' if cell_bold else '✗'}, "
            f"center={'✓' if cell_center else '✗'}"
        )

        bold_ok &= cell_bold
        center_ok &= cell_center

    # Scoring
    if bold_ok:
        print("✓ All header cells are bold (0.5 points)")
        score += 0.5
    else:
        print("✗ Not all header cells are bold (0 points)")

    if center_ok:
        print("✓ All header cells are centered (0.5 points)")
        score += 0.5
    else:
        print("✗ Not all header cells are centered (0 points)")

    final_score = min(score, max_score)
    print(f"Final Score: {final_score}")
    return final_score


if __name__ == "__main__":
    # Path provided by the task context
    pptx_path = "/home/user/set_the_header_row_of_table_1_to_bold_and_centered_text.pptx"
    reward = verify_header_bold_center(pptx_path)
    print(f"REWARD: {reward}")
