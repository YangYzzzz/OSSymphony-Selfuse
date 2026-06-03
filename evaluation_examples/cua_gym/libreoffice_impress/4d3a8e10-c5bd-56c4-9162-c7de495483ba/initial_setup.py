"""
Initial Setup: Algebra lesson presentation with 7 slides, slide 5 empty except for title.
Task ID: impress_teach_051
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_051'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_content(slide, title_text, body_lines):
    """Helper to set title and body content on a Title+Content slide."""
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide (layout 0) ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Algebra Lesson"
    slide1.placeholders[1].text = "Grade 8 Mathematics — Mrs. Patterson"

    # --- Slide 2: Introduction to Algebra (layout 1 = Title + Content) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_content(slide2, "Introduction to Algebra", [
        "Algebra is a branch of mathematics dealing with symbols and rules",
        "Variables represent unknown quantities we want to find",
        "Expressions combine numbers, variables, and operations",
        "Equations state that two expressions are equal",
        "Algebra is foundational for advanced math and science",
    ])

    # --- Slide 3: Variables and Constants (layout 1) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_content(slide3, "Variables and Constants", [
        "A variable is a letter that stands for a number (e.g., x, y, z)",
        "A constant is a fixed value (e.g., 3, -7, 0.5)",
        "Coefficients are numbers multiplied by variables (e.g., 2x means 2 times x)",
        "Terms are parts of an expression separated by + or -",
        "Example: In 3x + 7, the coefficient is 3 and the constant is 7",
    ])

    # --- Slide 4: Types of Equations (layout 1) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_content(slide4, "Types of Equations", [
        "Linear equations: highest power of variable is 1 (e.g., 2x + 5 = 15)",
        "Quadratic equations: highest power is 2 (e.g., x\u00b2 + 3x = 10)",
        "Systems of equations: two or more equations with shared variables",
        "We solve by isolating the variable on one side",
        "Always check your answer by substituting back",
    ])

    # --- Slide 5: Solving Linear Equations (layout 6 = Title Only) ---
    # This slide must have only the title and be otherwise EMPTY
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add just a title text box manually at the top
    from pptx.util import Emu
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Solving Linear Equations"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True

    # --- Slide 6: Practice Problems (layout 1) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_content(slide6, "Practice Problems", [
        "Solve each equation for x:",
        "1)  x + 9 = 21",
        "2)  3x - 4 = 14",
        "3)  5x + 2 = 27",
        "4)  4x - 10 = 18",
        "5)  7 + 2x = 23",
    ])

    # --- Slide 7: Summary & Homework (layout 1) ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_content(slide7, "Summary & Homework", [
        "Today we learned to solve one-step and two-step linear equations",
        "Remember: whatever you do to one side, do to the other",
        "Homework: Textbook page 142, exercises 1-20 (odd numbers)",
        "Quiz on solving linear equations next Wednesday",
        "Office hours: Tuesday and Thursday, 3:00 - 4:00 PM",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
