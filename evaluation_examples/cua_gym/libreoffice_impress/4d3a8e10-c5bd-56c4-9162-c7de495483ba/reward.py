"""
Reward Script: Animated step-by-step math problem on slide 5
Task ID: impress_teach_051
Domain: libreoffice_impress
Scoring:
  Component 1: Four text boxes with correct step text on slide 5 (0.50)
  Component 2: Vertical stacking order of text boxes (0.15)
  Component 3: Four 'Appear' animations present (0.20)
  Component 4: Animations are sequential on-click (0.15)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_051'

# Expected step texts (core content, checked with 'in' for flexibility)
EXPECTED_STEPS = [
    '2x + 5 = 15',
    '2x = 10',
    'x = 5',
    '2(5) + 5 = 15',
]

# More specific expected texts for exact matching
EXPECTED_FULL_TEXTS = [
    '1. Start with 2x + 5 = 15',
    '2. Subtract 5: 2x = 10',
    '3. Divide by 2: x = 5',
    '4. Check: 2(5) + 5 = 15',
]


def persist_app_state(domain):
    """Try to save any unsaved LibreOffice changes."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print(f"PERSIST: ctrl+s sent for {domain}")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide5 = prs.slides[4]

    # =========================================================================
    # Component 1: Four text boxes with correct step text on slide 5 (0.50 pts)
    # =========================================================================
    try:
        # Collect all text boxes on slide 5 (exclude the title placeholder and
        # the existing "Solving Linear Equations" text box)
        step_textboxes = []
        for shape in slide5.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                # Skip empty shapes and the existing title/header
                if not text:
                    continue
                if 'solving linear equations' in text.lower():
                    continue
                step_textboxes.append((shape, text))

        matched_steps = 0
        for expected_key in EXPECTED_STEPS:
            for shape, text in step_textboxes:
                if expected_key in text:
                    matched_steps += 1
                    break

        if matched_steps == 4:
            print(f"PASS: Component 1 -- All 4 step text boxes found (0.50 pts)")
            total_score += 0.50
        elif matched_steps > 0:
            partial = round(0.50 * (matched_steps / 4), 2)
            print(f"PARTIAL: Component 1 -- {matched_steps}/4 step texts found ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- No step text boxes found. Texts on slide 5: {[t for _, t in step_textboxes]}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # =========================================================================
    # Component 2: Text boxes are stacked vertically (0.15 pts)
    # Each subsequent step should have a larger 'top' value than the previous
    # =========================================================================
    try:
        # Re-collect step textboxes with their top positions, ordered by match
        step_positions = []
        for expected_key in EXPECTED_STEPS:
            for shape, text in step_textboxes:
                if expected_key in text:
                    step_positions.append((shape.top, text))
                    break

        if len(step_positions) >= 4:
            tops = [pos[0] for pos in step_positions]
            is_stacked = all(tops[i] < tops[i+1] for i in range(len(tops)-1))
            if is_stacked:
                print(f"PASS: Component 2 -- Text boxes stacked vertically (top values: {tops}) (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 2 -- Text boxes not in vertical order (top values: {tops})")
        else:
            print(f"FAIL: Component 2 -- Not enough step text boxes to check vertical stacking ({len(step_positions)}/4)")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # =========================================================================
    # Component 3: Four 'Appear' animations present on slide 5 (0.20 pts)
    # presetID="1" and presetClass="entr" = Appear entrance animation
    # =========================================================================
    try:
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide5.xml') as f:
                root = ET.parse(f).getroot()
                timing = root.find('.//p:timing', ns)

                if timing is None:
                    print("FAIL: Component 3 -- No timing/animation element found on slide 5")
                else:
                    # Find all animation nodes with presetID="1" (Appear) and presetClass="entr"
                    appear_anims = []
                    for ctn in timing.iter():
                        if ctn.tag.endswith('}cTn'):
                            preset_id = ctn.get('presetID')
                            preset_class = ctn.get('presetClass')
                            if preset_id == '1' and preset_class == 'entr':
                                appear_anims.append(ctn)

                    if len(appear_anims) >= 4:
                        print(f"PASS: Component 3 -- Found {len(appear_anims)} Appear animations (0.20 pts)")
                        total_score += 0.20
                    elif len(appear_anims) > 0:
                        partial = round(0.20 * (len(appear_anims) / 4), 2)
                        print(f"PARTIAL: Component 3 -- Found {len(appear_anims)}/4 Appear animations ({partial} pts)")
                        total_score += partial
                    else:
                        print(f"FAIL: Component 3 -- No Appear animations found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # =========================================================================
    # Component 4: Animations are sequential on-click (0.15 pts)
    # Each animation should be in a separate <p:par> with nodeType="clickEffect"
    # within the mainSeq, meaning each triggers on a separate click
    # =========================================================================
    try:
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide5.xml') as f:
                root = ET.parse(f).getroot()
                timing = root.find('.//p:timing', ns)

                if timing is None:
                    print("FAIL: Component 4 -- No timing element for sequence check")
                else:
                    # Count clickEffect nodes
                    click_effects = []
                    for ctn in timing.iter():
                        if ctn.tag.endswith('}cTn'):
                            if ctn.get('nodeType') == 'clickEffect':
                                click_effects.append(ctn)

                    # Also verify they target different shapes (sequential)
                    targeted_spids = set()
                    for ctn in timing.iter():
                        if ctn.tag.endswith('}spTgt'):
                            spid = ctn.get('spid')
                            if spid:
                                targeted_spids.add(spid)

                    if len(click_effects) >= 4 and len(targeted_spids) >= 4:
                        print(f"PASS: Component 4 -- {len(click_effects)} click-triggered animations targeting {len(targeted_spids)} shapes (0.15 pts)")
                        total_score += 0.15
                    else:
                        print(f"FAIL: Component 4 -- clickEffects={len(click_effects)}, unique targets={len(targeted_spids)} (need >=4 each)")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
