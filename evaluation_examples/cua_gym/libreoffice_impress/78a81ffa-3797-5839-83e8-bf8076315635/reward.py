"""
FINAL REWARD SCRIPT - SUCCESS
Task: Shade every other row (banded rows) in Table 1 with 10% gray.
Generated: 2025-10-17 13:34:14
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation
from pptx.enum.dml import MSO_FILL, MSO_COLOR_TYPE
from pptx.dml.color import RGBColor

"""
Reward Script
Task: "Shade every other row (banded rows) in Table 1 with 10% gray."
This script verifies that Table 1 (first table found in the presentation)
  • has every other row shaded (starting with the 2nd row)
  • shaded rows use 10 % gray (RGB≈230,230,230)
The script awards progressive points:
  60 %  – Correct banding pattern (row alternation)
  40 %  – Correct 10 %-gray colour on all shaded rows
It prints detailed diagnostics and finally prints
    REWARD: <score_between_0_and_1>
Exact 1.0 is given only when both pattern and colour are fully correct.
"""

FILE_PATH = "/home/user/shade_every_other_row_banded_rows_in_table_1_with_10_gray.pptx"

# -------------  helper utilities -------------

def _rgb_tuple(rgb):
    """Return (r,g,b) tuple from various python-pptx rgb representations."""
    if rgb is None:
        return None
    if isinstance(rgb, bytes):
        return (rgb[0], rgb[1], rgb[2])
    if isinstance(rgb, int):
        return ((rgb >> 16) & 0xFF, (rgb >> 8) & 0xFF, rgb & 0xFF)
    if isinstance(rgb, RGBColor):
        return (rgb[0], rgb[1], rgb[2])
    return None

def _cell_rgb(cell):
    """Return RGB tuple if the cell has a solid fill with an explicit RGB colour."""
    fill = cell.fill
    if fill and fill.type == MSO_FILL.SOLID:
        fc = fill.fore_color
        if fc.type == MSO_COLOR_TYPE.RGB and fc.rgb is not None:
            return _rgb_tuple(fc.rgb)
    return None

# -------------  main verification -------------

def verify_banded_rows(file_path: str) -> float:
    # Preliminary file checks
    if not os.path.exists(file_path):
        print("✗ File not found:", file_path)
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print("✗ Could not open PPTX:", e)
        return 0.0

    # Locate first table (Table 1)
    table = None
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                table = shape.table
                break
        if table is not None:
            break

    if table is None:
        print("✗ No table found in the presentation")
        return 0.0

    rows = table.rows
    total_rows = len(rows)
    if total_rows == 0:
        print("✗ Table contains no rows")
        return 0.0

    print(f"✓ Found table with {total_rows} rows and {len(table.columns)} columns")

    EXPECTED_GRAY = (230, 230, 230)  # 10 % gray in PowerPoint
    tolerance = 5                    # per-channel tolerance

    correctly_banded_rows = 0        # rows whose shaded/un-shaded status is correct
    correctly_coloured_rows = 0      # shaded rows that also have correct colour

    for idx, row in enumerate(rows):
        expected_shaded = (idx % 2 == 1)    # every 2nd row (starting with row 1)

        shaded_cells  = 0   # cells that have some solid fill
        gray_cells    = 0   # cells whose fill matches expected gray
        for cell in row.cells:
            rgb = _cell_rgb(cell)
            if rgb is not None:
                shaded_cells += 1
                if all(abs(rgb[i] - EXPECTED_GRAY[i]) <= tolerance for i in range(3)):
                    gray_cells += 1

        # A row counts as *shaded* if the majority of its cells are shaded
        is_row_shaded = shaded_cells >= (len(row.cells) / 2)
        is_gray_correct = (shaded_cells > 0 and gray_cells == shaded_cells)

        print(
            f"Row {idx}: shaded_cells={shaded_cells}/{len(row.cells)}, "
            f"is_row_shaded={is_row_shaded}, is_gray_correct={is_gray_correct}, "
            f"expected_shaded={expected_shaded}"
        )

        # Banding pattern correctness
        if is_row_shaded == expected_shaded:
            correctly_banded_rows += 1

        # Colour correctness only matters for rows that should be shaded
        if expected_shaded and is_row_shaded and is_gray_correct:
            correctly_coloured_rows += 1

    # ---------- scoring ----------
    banding_score = correctly_banded_rows / total_rows
    expected_shaded_rows = total_rows // 2   # floor division adequate for scoring
    color_score = (correctly_coloured_rows / expected_shaded_rows) if expected_shaded_rows else 0.0

    # Weighted total: 60 % banding pattern, 40 % colour accuracy
    total_score = 0.6 * banding_score + 0.4 * color_score
    total_score = max(0.0, min(1.0, total_score))  # clamp

    print(f"Banding pattern correctness: {banding_score:.2f} (weight 0.6)")
    print(f"Shade colour correctness:    {color_score:.2f} (weight 0.4)")
    print(f"Total verification score:    {total_score:.2f}")

    return total_score

# -------------  script entry point -------------

if __name__ == "__main__":
    reward = verify_banded_rows(FILE_PATH)
    print(f"REWARD: {reward}")

