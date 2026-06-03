"""
Initial Setup: Add a rounded rectangle callout box on slide 5
Task ID: impress_teach_032
Domain: libreoffice_impress

Creates a 7-slide Research Methods presentation. Slide 5 has content about
citations but NO callout box (that is the task for the agent to add).
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_032'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color palette
    TITLE_BG = RGBColor(0x1B, 0x3A, 0x5C)   # Dark navy
    CONTENT_BG = RGBColor(0xF5, 0xF5, 0xF5)  # Light grey
    ACCENT = RGBColor(0x2E, 0x86, 0xC1)       # Blue accent
    DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    def set_bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title_shape(slide, text, left, top, width, height, font_size, color, bold=True):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        return txBox

    def add_bullet_content(slide, items, left, top, width, height, font_size=16, color=DARK_TEXT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.space_after = Pt(6)
            run = p.runs[0]
            run.font.name = "Arial"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color

    # ===== Slide 1: Title Slide =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_bg(slide1, TITLE_BG)
    add_title_shape(slide1, "Research Methods", Inches(1), Inches(2), Inches(10), Inches(1.5), 44, WHITE)
    add_title_shape(slide1, "A Comprehensive Guide for Graduate Students", Inches(1), Inches(3.5), Inches(10), Inches(1), 22, RGBColor(0xAE, 0xD6, 0xF1), bold=False)
    add_title_shape(slide1, "Dr. Elena Vasquez  |  Department of Social Sciences  |  Spring 2026", Inches(1), Inches(5.5), Inches(10), Inches(0.6), 14, RGBColor(0x85, 0xC1, 0xE9), bold=False)

    # ===== Slide 2: Types of Research =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide2, CONTENT_BG)
    add_title_shape(slide2, "Types of Research", Inches(0.8), Inches(0.4), Inches(10), Inches(0.8), 32, ACCENT)
    add_bullet_content(slide2, [
        "Quantitative Research: Uses numerical data and statistical analysis",
        "Qualitative Research: Explores meanings, experiences, and social phenomena",
        "Mixed Methods: Combines quantitative and qualitative approaches",
        "Experimental: Manipulates variables to establish causal relationships",
        "Observational: Collects data through systematic observation",
        "Survey-based: Gathers information using structured questionnaires",
    ], Inches(1), Inches(1.5), Inches(10), Inches(4.5))

    # ===== Slide 3: Literature Review =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide3, CONTENT_BG)
    add_title_shape(slide3, "Literature Review Process", Inches(0.8), Inches(0.4), Inches(10), Inches(0.8), 32, ACCENT)
    add_bullet_content(slide3, [
        "1. Define your research question and scope",
        "2. Search academic databases (JSTOR, PubMed, Google Scholar)",
        "3. Screen sources for relevance and quality",
        "4. Extract key findings and methodological details",
        "5. Synthesize themes across multiple studies",
        "6. Identify gaps in existing knowledge",
        "7. Write a critical narrative connecting sources to your work",
    ], Inches(1), Inches(1.5), Inches(10), Inches(4.5))

    # ===== Slide 4: Data Collection =====
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide4, CONTENT_BG)
    add_title_shape(slide4, "Data Collection Methods", Inches(0.8), Inches(0.4), Inches(10), Inches(0.8), 32, ACCENT)
    add_bullet_content(slide4, [
        "Interviews: Semi-structured conversations with participants",
        "Surveys: Online or paper questionnaires (n > 100 recommended)",
        "Focus Groups: Facilitated group discussions (6-10 participants)",
        "Case Studies: In-depth examination of a single entity or event",
        "Archival Research: Analysis of existing records and documents",
        "Ethnography: Immersive fieldwork within a community or culture",
    ], Inches(1), Inches(1.5), Inches(10), Inches(4.5))

    # ===== Slide 5: Citing Sources (NO callout - that is the task) =====
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide5, CONTENT_BG)
    add_title_shape(slide5, "Citing Sources in Academic Work", Inches(0.8), Inches(0.4), Inches(10), Inches(0.8), 32, ACCENT)
    add_bullet_content(slide5, [
        "APA Style: Author-date format used widely in social sciences",
        "   Example: (Vasquez, 2024) or Vasquez (2024) argued that...",
        "MLA Style: Author-page format common in humanities",
        "   Example: (Vasquez 42) with a Works Cited page",
        "Chicago Style: Notes-bibliography or author-date variants",
        "   Example: Footnotes with full bibliographic information",
        "IEEE Style: Numbered references in brackets for engineering",
        "   Example: [1] with a numbered reference list",
    ], Inches(1), Inches(1.5), Inches(10), Inches(4.5))

    # ===== Slide 6: Analysis Methods =====
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide6, CONTENT_BG)
    add_title_shape(slide6, "Analysis Methods", Inches(0.8), Inches(0.4), Inches(10), Inches(0.8), 32, ACCENT)
    add_bullet_content(slide6, [
        "Descriptive Statistics: Mean, median, mode, standard deviation",
        "Inferential Statistics: t-tests, ANOVA, regression analysis",
        "Thematic Analysis: Identifying patterns and themes in qualitative data",
        "Content Analysis: Systematic coding of textual or media content",
        "Grounded Theory: Building theory from iterative data analysis",
        "Discourse Analysis: Examining language use in social contexts",
    ], Inches(1), Inches(1.5), Inches(10), Inches(4.5))

    # ===== Slide 7: Conclusion =====
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide7, TITLE_BG)
    add_title_shape(slide7, "Key Takeaways", Inches(1), Inches(1), Inches(10), Inches(1), 36, WHITE)
    add_bullet_content(slide7, [
        "Choose methods aligned with your research questions",
        "Maintain ethical standards throughout data collection",
        "Document every decision for reproducibility",
        "Seek peer feedback at every stage of your research",
        "Stay current with evolving methodological best practices",
    ], Inches(1), Inches(2.3), Inches(10), Inches(4), font_size=18, color=WHITE)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
