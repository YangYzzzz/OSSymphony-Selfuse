"""
Reward Script: Add rounded rectangle callout on slide 5
Task ID: impress_teach_032
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.35): Rounded rectangle shape on slide 5 with correct text
  - Component 2 (0.25): Fill color #FFF9C4
  - Component 3 (0.25): Border color #F9A825 with 1.5pt width
  - Component 4 (0.15): Shape positioned in bottom-right corner
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_032'
FILE_PATH = os.path.join(WORKDIR, f'{TASK_ID}.pptx')


def persist_app_state(domain):
    """Save any unsaved changes in LibreOffice."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Missing library: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"PRECONDITION FAIL: Need at least 5 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # Slide 5 (0-indexed)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Find the callout shape: look for a shape with the target text that is NOT a plain textbox
    # from initial. We look for any shape containing the target text.
    TARGET_TEXT = "Remember: Always cite your sources!"
    callout_shape = None

    for shape in slide.shapes:
        if shape.has_text_frame:
            shape_text = ""
            for para in shape.text_frame.paragraphs:
                shape_text += para.text
            if TARGET_TEXT.lower() in shape_text.strip().lower():
                callout_shape = shape
                break

    # Component 1: Rounded rectangle shape on slide 5 with correct text (0.35 points)
    try:
        if callout_shape is not None:
            # Verify it's an auto shape (rounded rectangle), not just a textbox
            is_auto_shape = (callout_shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE)
            # Also accept FREEFORM or other shape types that could be a rounded rectangle
            # The key is it has the text and is not one of the original textboxes
            shape_text_full = ""
            for para in callout_shape.text_frame.paragraphs:
                shape_text_full += para.text
            text_matches = TARGET_TEXT.lower() in shape_text_full.strip().lower()

            if text_matches and is_auto_shape:
                print(f"PASS: Component 1 - Rounded rectangle with correct text found (0.35 pts)")
                print(f"  Shape type: {callout_shape.shape_type}, name: {callout_shape.name}")
                total_score += 0.35
            elif text_matches:
                # Text matches but not an auto shape - partial credit
                print(f"PARTIAL: Component 1 - Text found but shape type is {callout_shape.shape_type}, not AUTO_SHAPE (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 1 - Text does not match. Found: {repr(shape_text_full)}")
        else:
            print(f"FAIL: Component 1 - No shape with text '{TARGET_TEXT}' found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Components 2-4 only checked if callout shape was found
    if callout_shape is None:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Fill color is #FFF9C4 (0.25 points)
    try:
        fill = callout_shape.fill
        if fill.type is not None and fill.type == 1:  # SOLID fill
            fill_rgb = str(fill.fore_color.rgb).upper()
            if fill_rgb == "FFF9C4":
                print(f"PASS: Component 2 - Fill color is #FFF9C4 (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 2 - Fill color is #{fill_rgb}, expected #FFF9C4")
        else:
            print(f"FAIL: Component 2 - No solid fill found (fill type: {fill.type})")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Border color #F9A825 with 1.5pt width (0.25 points)
    try:
        line = callout_shape.line
        comp3_score = 0.0

        # Sub-check A: line color == #F9A825
        try:
            detected_color = None
            if line.color and line.color.type is not None:
                detected_color = str(line.color.rgb).upper()
            elif line.fill.type == 1:  # SOLID fill fallback
                detected_color = str(line.fill.fore_color.rgb).upper()

            if detected_color == "F9A825":
                print(f"  Border color: #{detected_color} (correct)")
                comp3_score += 0.12
            else:
                print(f"  Border color: #{detected_color} (expected #F9A825)")
        except Exception as e:
            print(f"  Border color check error: {e}")

        # Sub-check B: line width == 1.5pt (19050 EMU)
        try:
            expected_width = 19050  # 1.5pt in EMU
            actual_width = line.width
            if actual_width is not None and abs(actual_width - expected_width) / expected_width <= 0.1:
                print(f"  Border width: {actual_width} EMU = {actual_width/12700:.2f}pt (correct)")
                comp3_score += 0.13
            else:
                print(f"  Border width: {actual_width} EMU (expected 19050 EMU / 1.5pt)")
        except Exception as e:
            print(f"  Border width check error: {e}")

        if comp3_score >= 0.24:
            print(f"PASS: Component 3 - Border color and width correct (0.25 pts)")
            total_score += 0.25
        elif comp3_score > 0:
            print(f"PARTIAL: Component 3 - ({comp3_score:.2f} pts)")
            total_score += comp3_score
        else:
            print(f"FAIL: Component 3 - Neither border color nor width correct")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Shape positioned in bottom-right corner (0.15 points)
    try:
        left = callout_shape.left
        top = callout_shape.top
        width = callout_shape.width
        height = callout_shape.height

        right_edge = left + width
        bottom_edge = top + height

        # "Bottom-right corner" means:
        # - The right edge of the shape should be in the right half of the slide
        # - The bottom edge should be in the bottom half of the slide
        # More specifically, the shape center should be in the bottom-right quadrant
        center_x = left + width / 2
        center_y = top + height / 2

        in_right_half = center_x > slide_width / 2
        in_bottom_half = center_y > slide_height / 2

        print(f"  Shape center: ({center_x}, {center_y})")
        print(f"  Slide midpoint: ({slide_width/2}, {slide_height/2})")
        print(f"  Right half: {in_right_half}, Bottom half: {in_bottom_half}")

        if in_right_half and in_bottom_half:
            print(f"PASS: Component 4 - Shape is in bottom-right corner (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 - Shape not in bottom-right corner")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(FILE_PATH)
