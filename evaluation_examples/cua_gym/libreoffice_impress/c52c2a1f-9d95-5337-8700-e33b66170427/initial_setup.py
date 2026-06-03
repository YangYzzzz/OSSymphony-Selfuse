"""
Initial Setup: Science Fair presentation with 8 slides. Slide 5 'Results' has no transition.
Task ID: impress_tm_018
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_018'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.paragraphs[0].text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
    return slide


def add_blank_with_textbox(prs, title_text, body_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    # Title textbox
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.0), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
    p.alignment = PP_ALIGN.LEFT
    # Body textbox
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.0), Inches(4.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = body_text
    p2.font.size = Pt(16)
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    add_title_slide(
        prs,
        "Solar Energy Harvesting Efficiency",
        "Regional Science Fair 2025 — Westfield Academy\nPresented by Ava Martinez & Liam Nakamura"
    )

    # --- Slide 2: Introduction ---
    add_content_slide(prs, "Introduction", [
        "Photovoltaic cell efficiency varies significantly with panel angle",
        "Current residential installations use fixed 30-degree mounting",
        "Dynamic tracking systems cost 3x more but promise 25-40% gain",
        "Research question: Can a low-cost adjustable mount match tracker output?",
        "Hypothesis: Manual bi-weekly adjustment achieves at least 85% of tracker yield",
    ])

    # --- Slide 3: Materials & Methods ---
    add_content_slide(prs, "Materials & Methods", [
        "Three identical 50W monocrystalline panels (SunPower SPR-E20)",
        "Panel A: Fixed at 30 degrees (control)",
        "Panel B: Manual tilt adjustment every 14 days",
        "Panel C: Motorized single-axis tracker (reference)",
        "Data logged every 15 minutes via Arduino Mega + INA219 sensors",
        "Test period: March 1 – August 31, 2025 (184 days)",
        "Location: 37.7749 N, 122.4194 W (San Francisco, CA)",
    ])

    # --- Slide 4: Data Collection ---
    add_content_slide(prs, "Data Collection Overview", [
        "Total data points collected: 176,256 (3 panels x 96/day x 184 days)",
        "Weather data sourced from NOAA station USW00023234",
        "Cloud cover categorized: Clear (<25%), Partial (25-75%), Overcast (>75%)",
        "Panel temperature monitored with DS18B20 waterproof probes",
        "Inverter efficiency normalized across all three setups",
        "Outliers removed using IQR method (< 0.3% of readings)",
    ])

    # --- Slide 5: Results (NO transition — this is the task target) ---
    add_content_slide(prs, "Results", [
        "Panel A (fixed): 38.2 kWh/month average",
        "Panel B (manual adjust): 46.7 kWh/month average (+22.3%)",
        "Panel C (tracker): 51.4 kWh/month average (+34.6%)",
        "Manual adjustment captured 90.9% of tracker performance",
        "Largest gains observed in March and September (equinox periods)",
        "Cloud cover reduced differential between methods by 40-60%",
        "Peak single-day output: Panel C at 2.84 kWh on June 21",
    ])

    # --- Slide 6: Statistical Analysis ---
    add_content_slide(prs, "Statistical Analysis", [
        "One-way ANOVA: F(2, 549) = 47.32, p < 0.001",
        "Tukey HSD post-hoc: All pairwise comparisons significant (p < 0.01)",
        "Effect size (Cohen's d): A vs B = 1.12, A vs C = 1.89, B vs C = 0.64",
        "Linear regression R² = 0.87 for angle-vs-output model",
        "Seasonal adjustment coefficient: 0.93 (spring) to 1.07 (summer)",
    ])

    # --- Slide 7: Conclusions ---
    add_content_slide(prs, "Conclusions", [
        "Hypothesis supported: manual adjustment achieved 90.9% of tracker yield",
        "Cost-benefit ratio favors manual adjustment for residential use",
        "Estimated ROI: Manual mount pays back in 2.1 years vs 5.7 for tracker",
        "Bi-weekly adjustment requires only 10 minutes per session",
        "Recommended optimal angles published in supplementary table",
    ])

    # --- Slide 8: References & Acknowledgments ---
    add_content_slide(prs, "References & Acknowledgments", [
        "Green, M.A. et al. (2024). Solar cell efficiency tables. Prog. Photovolt. 32(1)",
        "NREL. (2024). Best Research-Cell Efficiency Chart. Golden, CO",
        "Kalogirou, S. (2023). Solar Energy Engineering, 3rd ed. Academic Press",
        "Thanks to Dr. Patricia Hwang (Westfield Academy Physics Dept.)",
        "Special thanks to Bay Area Solar Cooperative for equipment loans",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
