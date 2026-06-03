"""
Initial Setup: Thesis Defense presentation and feedback document for notes import task
Task ID: osworld_multi_apps_impress_notes_import_015
Domain: libreoffice_impress (multi-app: also creates feedback.docx)
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from docx import Document

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_impress_notes_import_015'
PPTX_PATH = f'{WORKDIR}/Thesis_Defense.pptx'
DOCX_PATH = f'{DESKTOP}/feedback.docx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


# Slide definitions: (title, body_bullets)
SLIDES = [
    (
        "Thesis Defense: Adaptive Machine Learning for Climate Prediction",
        [
            "Doctoral Candidate: Emily Zhang",
            "Department of Atmospheric and Computational Sciences",
            "Advisor: Prof. Michael Hartwell",
            "Committee Members: Dr. Sara Okonkwo, Dr. Yusuf Patel, Dr. Linda Brosseau",
        ],
    ),
    (
        "Research Motivation",
        [
            "Climate change accelerates extreme weather events globally",
            "Traditional numerical weather prediction (NWP) models are computationally expensive",
            "Data-driven approaches offer 100x speedup with comparable accuracy",
            "Gap: Existing ML models lack adaptability to distribution shift",
            "This research proposes AdaptCast: a continual-learning framework",
        ],
    ),
    (
        "Research Objectives",
        [
            "Objective 1: Develop adaptive neural architecture for temporal distribution shift",
            "Objective 2: Benchmark against ECMWF ERA5 reanalysis dataset",
            "Objective 3: Demonstrate zero-shot generalization to unseen climate zones",
            "Objective 4: Validate on real-time NOAA observation feeds",
        ],
    ),
    (
        "Literature Review",
        [
            "Reichstein et al. (2019): Deep learning in Earth system sciences",
            "Rolnick et al. (2022): Tackling climate change with machine learning",
            "Weyn et al. (2020): Improving data-driven global weather prediction",
            "Nguyen et al. (2023): ClimaX — foundation model for weather and climate",
            "Key gap: None address continual learning under concept drift",
        ],
    ),
    (
        "Methodology Overview",
        [
            "Phase 1: Data curation — ERA5, CMIP6, station observations (1979–2023)",
            "Phase 2: Baseline models — U-Net, GraphCast, Pangu-Weather",
            "Phase 3: AdaptCast architecture — meta-learning + elastic weight consolidation",
            "Phase 4: Evaluation — skill scores, RMSE, anomaly correlation coefficient",
        ],
    ),
    (
        "AdaptCast Architecture",
        [
            "Backbone: Vision Transformer (ViT-L/16) pretrained on ERA5",
            "Adapter modules: lightweight bottleneck layers per climate zone",
            "Continual learning: Elastic Weight Consolidation (EWC) regularizer",
            "Memory replay: Stratified reservoir sampling across decades",
            "Output heads: 10-day forecast, seasonal anomaly, extreme event probability",
        ],
    ),
    (
        "Dataset and Experimental Setup",
        [
            "ERA5 reanalysis: 0.25° × 0.25° global grid, 37 pressure levels",
            "Training: 1979–2010 (32 years, ~2.4 TB)",
            "Validation: 2011–2015 (5 years)",
            "Test: 2016–2023 (8 years, includes COVID-era anomalies)",
            "Hardware: 64× A100 GPU cluster, 6 weeks total training",
        ],
    ),
    (
        "Results: Forecast Accuracy",
        [
            "Z500 RMSE at 5-day: AdaptCast 312 m²/s² vs GraphCast 334 m²/s² (6.6% improvement)",
            "T850 bias: AdaptCast reduces warm bias from +0.8 K to +0.2 K",
            "Tropical cyclone track error: 15% reduction at 72-hour lead time",
            "Extreme precipitation events: F1-score 0.71 vs baseline 0.58",
        ],
    ),
    (
        "Results: Continual Learning Evaluation",
        [
            "Catastrophic forgetting reduction: 41% less accuracy degradation vs fine-tuning",
            "EWC λ=5000 optimal across all climate zones",
            "Zero-shot performance on Antarctic domain: 78% of supervised performance",
            "Adapter parameter overhead: only 2.3% of total model parameters",
        ],
    ),
    (
        "Ablation Study",
        [
            "Without EWC: 23% higher forgetting on historical benchmarks",
            "Without adapters: 11% worse on domain-specific extremes",
            "Without memory replay: oscillating validation loss, fails to converge",
            "Full model (AdaptCast): best across all 14 evaluation metrics",
        ],
    ),
    (
        "Discussion and Limitations",
        [
            "Strength: AdaptCast generalizes across climate regimes without full retraining",
            "Strength: EWC overhead negligible at inference time",
            "Limitation: Data-hungry — requires >10 years historical data per zone",
            "Limitation: Cannot yet assimilate real-time satellite swath data",
            "Future work: Incorporate GOES-18 imagery and ocean heat content",
        ],
    ),
    (
        "Conclusion and Future Directions",
        [
            "AdaptCast achieves state-of-the-art climate forecasting with continual learning",
            "Demonstrated robust generalization under temporal distribution shift",
            "Open-source release: github.com/ezhang-lab/adaptcast",
            "Next steps: Operational trial with NOAA Environmental Modeling Center",
            "Thank you — Questions welcome",
        ],
    ),
]

# Feedback notes content: realistic presenter notes for each slide
NOTES_CONTENT = [
    # Slide 1 - Title
    "Welcome the committee and thank them for their time. Mention this defense represents four years of doctoral research. Briefly state the problem domain before moving to motivation.",

    # Slide 2 - Motivation
    "Emphasize the computational cost gap: NWP models take 10–30 minutes per forecast cycle even on supercomputers. ML models run in seconds. Highlight that distribution shift is the central unsolved challenge that motivates this entire thesis.",

    # Slide 3 - Objectives
    "Walk through each objective in order. Note that Objective 3 (zero-shot) is particularly novel and has not been attempted in prior work. Expect committee questions about how zero-shot is defined and evaluated.",

    # Slide 4 - Literature Review
    "Acknowledge Nguyen et al. ClimaX as the closest prior work. Clarify how AdaptCast differs: ClimaX is a static pretrained model, whereas AdaptCast actively adapts post-deployment. This is the key scientific contribution.",

    # Slide 5 - Methodology
    "Spend about two minutes on this slide. Committees often want to understand the overall pipeline before diving into architecture details. Point out that Phase 1 data curation took six months due to quality control issues in station data.",

    # Slide 6 - Architecture
    "This is a technical audience, so explain EWC briefly: it adds a quadratic penalty to the loss function, weighted by Fisher information, to prevent overwriting important weights. The adapter design is borrowed from NLP transfer learning literature.",

    # Slide 7 - Dataset and Setup
    "Acknowledge the large compute cost and mention this was made possible by a compute grant from NSF ACCESS. The COVID-era test period is important because it contains anomalous human activity signals that stress-test generalization.",

    # Slide 8 - Forecast Accuracy
    "Z500 is geopotential height at 500 hPa — a standard benchmark in NWP. The 6.6% RMSE improvement may sound modest but is highly significant at synoptic scale. Emphasize the tropical cyclone result as it has direct societal impact.",

    # Slide 9 - Continual Learning
    "The 41% reduction in catastrophic forgetting is the headline result of this thesis. Stress that EWC at λ=5000 was selected via grid search over 20 values. The 78% zero-shot result on Antarctica is particularly exciting given no Antarctic training data was used.",

    # Slide 10 - Ablation
    "Ablations confirm that each component contributes independently. The memory replay result is the most dramatic: without it, the model oscillates and never converges past epoch 30. This validates the reservoir sampling strategy.",

    # Slide 11 - Discussion
    "Be honest about limitations. Committees appreciate candor. The satellite assimilation limitation is real but addressed in the future work. Mention that a collaboration with NOAA is already in progress to tackle real-time data ingestion.",

    # Slide 12 - Conclusion
    "End with the open-source release — this signals impact beyond academia. Reiterate the three main contributions: architecture, continual learning method, and benchmark. Invite questions and thank committee again.",
]

# Content blocks (these go in feedback.docx [CONTENT] sections — different from notes)
CONTENT_REVISIONS = [
    "No content changes recommended for the title slide. Formatting is clean and professional.",
    "Consider adding a specific statistic: 'Global average temperature has risen 1.1°C since pre-industrial era (IPCC AR6, 2021).' Place after the first bullet.",
    "Objectives are clear. Consider renumbering as O1–O4 for easier reference during Q&A.",
    "Add Lam et al. (2023) GraphCast paper as it is directly benchmarked. Full citation: Lam, R. et al. (2023). Learning skillful medium-range global weather forecasting. Science, 382, 1416–1421.",
    "Rename 'Phase 1' to 'Stage 1' for consistency with the terminology used in Chapter 2 of the written thesis.",
    "Add a small architectural diagram reference: 'See Figure 3.2 in thesis for detailed layer dimensions.'",
    "Update hardware to reflect final allocation: '64× A100 80GB GPU cluster, 6 weeks total compute (equivalent to 2,688 GPU-days).'",
    "Add confidence intervals: 'RMSE improvement 6.6% ± 0.4% (95% CI, bootstrapped over 500 test windows).'",
    "Add statistical significance note: 'All improvements significant at p < 0.01 via Wilcoxon signed-rank test.'",
    "Add a comparison row to the ablation table: 'Elastic + Replay (no adapters)' to show additive effect.",
    "Soften the satellite limitation statement: 'Current version does not yet support real-time swath ingestion; this is addressed in ongoing NOAA collaboration (see Future Work).'",
    "Update GitHub URL to: github.com/ezhang-lab/adaptcast (v1.2.0 tagged at thesis submission). Add DOI: 10.5281/zenodo.9876543.",
]


def create_pptx():
    """Create the initial Thesis_Defense.pptx with 12 slides, NO notes."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]   # Title Slide
    content_layout = prs.slide_layouts[1]  # Title and Content

    for i, (title_text, bullets) in enumerate(SLIDES):
        if i == 0:
            slide = prs.slides.add_slide(title_layout)
            slide.shapes.title.text = title_text
            # Subtitle placeholder
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = "\n".join(bullets)
        else:
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = title_text
            tf = slide.placeholders[1].text_frame
            tf.text = bullets[0]
            for bullet in bullets[1:]:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
        # NOTE: Do NOT set any notes — the task requires the agent to insert notes

    prs.save(PPTX_PATH)
    print(f'Initial PPTX created: {PPTX_PATH}')


def create_feedback_docx():
    """Create feedback.docx on the Desktop with [CONTENT] and [NOTES] tagged sections."""
    os.makedirs(DESKTOP, exist_ok=True)

    doc = Document()
    doc.add_heading('Mentor Feedback: Thesis Defense Slides', level=0)

    intro = doc.add_paragraph(
        'Dear Emily, please find below my structured feedback for each slide. '
        'Each section contains a [CONTENT] block with revision suggestions for the slide body, '
        'and a [NOTES] block with suggested presenter notes you should add to the Notes panel in Impress. '
        'Please incorporate only the [NOTES] sections into the presentation file.'
    )

    for i, (title_text, _) in enumerate(SLIDES):
        slide_num = i + 1

        # Section heading
        doc.add_heading(f'Slide {slide_num}: {title_text}', level=1)

        # [CONTENT] block
        doc.add_paragraph('[CONTENT]', style='Intense Quote')
        doc.add_paragraph(CONTENT_REVISIONS[i])

        # [NOTES] block
        doc.add_paragraph('[NOTES]', style='Intense Quote')
        doc.add_paragraph(NOTES_CONTENT[i])

        # Separator (blank line)
        doc.add_paragraph('')

    doc.save(DOCX_PATH)
    print(f'Feedback DOCX created: {DOCX_PATH}')


def main():
    create_pptx()
    create_feedback_docx()

    # GUI-ready startup: open the presentation in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{PPTX_PATH}"', delay_sec=3.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


main()
