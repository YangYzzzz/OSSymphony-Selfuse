"""
Reward Script: Parse feedback.docx, extract [NOTES] sections, insert into Thesis_Defense.pptx slide notes
Task ID: osworld_multi_apps_impress_notes_import_015
Domain: libreoffice_impress
Scoring:
  Component 1: All 12 slides have non-empty notes text             (0.4 pts)
  Component 2: Notes content accuracy for slides 1-6               (0.3 pts)
  Component 3: Notes content accuracy for slides 7-12              (0.3 pts)
  Total: 1.0
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_multi_apps_impress_notes_import_015'

# Expected notes for each slide (derived from feedback.docx [NOTES] sections)
EXPECTED_NOTES = {
    1: "Welcome the committee and thank them for their time. Mention this defense represents four years of doctoral research. Briefly state the problem domain before moving to motivation.",
    2: "Emphasize the computational cost gap: NWP models take 10\u201330 minutes per forecast cycle even on supercomputers. ML models run in seconds. Highlight that distribution shift is the central unsolved challenge that motivates this entire thesis.",
    3: "Walk through each objective in order. Note that Objective 3 (zero-shot) is particularly novel and has not been attempted in prior work. Expect committee questions about how zero-shot is defined and evaluated.",
    4: "Acknowledge Nguyen et al. ClimaX as the closest prior work. Clarify how AdaptCast differs: ClimaX is a static pretrained model, whereas AdaptCast actively adapts post-deployment. This is the key scientific contribution.",
    5: "Spend about two minutes on this slide. Committees often want to understand the overall pipeline before diving into architecture details. Point out that Phase 1 data curation took six months due to quality control issues in station data.",
    6: "This is a technical audience, so explain EWC briefly: it adds a quadratic penalty to the loss function, weighted by Fisher information, to prevent overwriting important weights. The adapter design is borrowed from NLP transfer learning literature.",
    7: "Acknowledge the large compute cost and mention this was made possible by a compute grant from NSF ACCESS. The COVID-era test period is important because it contains anomalous human activity signals that stress-test generalization.",
    8: "Z500 is geopotential height at 500 hPa \u2014 a standard benchmark in NWP. The 6.6% RMSE improvement may sound modest but is highly significant at synoptic scale. Emphasize the tropical cyclone result as it has direct societal impact.",
    9: "The 41% reduction in catastrophic forgetting is the headline result of this thesis. Stress that EWC at \u03bb=5000 was selected via grid search over 20 values. The 78% zero-shot result on Antarctica is particularly exciting given no Antarctic training data was used.",
    10: "Ablations confirm that each component contributes independently. The memory replay result is the most dramatic: without it, the model oscillates and never converges past epoch 30. This validates the reservoir sampling strategy.",
    11: "Be honest about limitations. Committees appreciate candor. The satellite assimilation limitation is real but addressed in the future work. Mention that a collaboration with NOAA is already in progress to tackle real-time data ingestion.",
    12: "End with the open-source release \u2014 this signals impact beyond academia. Reiterate the three main contributions: architecture, continual learning method, and benchmark. Invite questions and thank committee again.",
}


def get_slide_notes(slide):
    """Safely get notes text from a slide."""
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def normalize_notes(text):
    """Normalize whitespace for comparison."""
    return " ".join(text.split())


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Sanity check: exactly 12 slides
    num_slides = len(prs.slides)
    if num_slides != 12:
        print(f"WARN: Expected 12 slides, found {num_slides}")

    # Component 1: All 12 slides have non-empty notes (0.4 points)
    # This FAILS on initial (all empty) → PASSES on golden (all populated)
    try:
        notes_per_slide = {}
        for i, slide in enumerate(prs.slides):
            notes = get_slide_notes(slide)
            notes_per_slide[i + 1] = notes

        slides_with_notes = sum(1 for n in notes_per_slide.values() if n)
        if slides_with_notes == 12:
            print(f"PASS: Component 1 — All 12 slides have non-empty notes (0.4 pts)")
            total_score += 0.4
        elif slides_with_notes > 0:
            partial = round(0.4 * slides_with_notes / 12, 4)
            print(f"PARTIAL: Component 1 — {slides_with_notes}/12 slides have non-empty notes ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — No slides have notes (expected all 12 to have notes)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Notes content accuracy for slides 1-6 (0.3 points)
    # Each slide contributes 0.05 points if notes match expected text
    # This FAILS on initial (empty notes) → PASSES on golden (notes match expected)
    try:
        correct_slides_1_6 = 0
        for slide_num in range(1, 7):
            if slide_num > len(prs.slides):
                break
            actual_notes = notes_per_slide.get(slide_num, "")
            expected = EXPECTED_NOTES.get(slide_num, "")
            actual_norm = normalize_notes(actual_notes)
            expected_norm = normalize_notes(expected)
            if actual_norm == expected_norm:
                correct_slides_1_6 += 1
                print(f"PASS: Component 2 — Slide {slide_num} notes match expected text")
            else:
                print(f"FAIL: Component 2 — Slide {slide_num} notes mismatch")
                print(f"  Expected (first 80 chars): {expected_norm[:80]!r}")
                print(f"  Actual   (first 80 chars): {actual_norm[:80]!r}")

        comp2_score = round(0.3 * correct_slides_1_6 / 6, 4)
        print(f"Component 2 subtotal: {correct_slides_1_6}/6 slides correct → {comp2_score} pts")
        if comp2_score > 0:
            total_score += comp2_score
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Notes content accuracy for slides 7-12 (0.3 points)
    # Each slide contributes 0.05 points if notes match expected text
    # This FAILS on initial (empty notes) → PASSES on golden (notes match expected)
    try:
        correct_slides_7_12 = 0
        for slide_num in range(7, 13):
            if slide_num > len(prs.slides):
                break
            actual_notes = notes_per_slide.get(slide_num, "")
            expected = EXPECTED_NOTES.get(slide_num, "")
            actual_norm = normalize_notes(actual_notes)
            expected_norm = normalize_notes(expected)
            if actual_norm == expected_norm:
                correct_slides_7_12 += 1
                print(f"PASS: Component 3 — Slide {slide_num} notes match expected text")
            else:
                print(f"FAIL: Component 3 — Slide {slide_num} notes mismatch")
                print(f"  Expected (first 80 chars): {expected_norm[:80]!r}")
                print(f"  Actual   (first 80 chars): {actual_norm[:80]!r}")

        comp3_score = round(0.3 * correct_slides_7_12 / 6, 4)
        print(f"Component 3 subtotal: {correct_slides_7_12}/6 slides correct → {comp3_score} pts")
        if comp3_score > 0:
            total_score += comp3_score
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 4)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/Thesis_Defense.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
