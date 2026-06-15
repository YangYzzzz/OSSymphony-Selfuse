"""
Reward Script: Animated process flow on slide 5 with four stages, arrows, descriptions, and animations
Task ID: impress_sales_094
Domain: libreoffice_impress
Scoring:
  C1 (0.25) — Four stage rounded rectangles with correct text labels
  C2 (0.20) — Stage shapes have correct fill colors
  C3 (0.15) — Three connecting arrows between stages
  C4 (0.15) — Four 2-line description text boxes below stages
  C5 (0.15) — Stage shapes have Zoom animation on click
  C6 (0.10) — Arrows have Wipe animation, descriptions have Fade animation
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_094'

# Expected stage labels and colors
EXPECTED_STAGES = ['Discover', 'Design', 'Deliver', 'Delight']
EXPECTED_COLORS = {
    'Discover': '2B6CB0',
    'Design': '4CAF50',
    'Deliver': 'FF6B35',
    'Delight': '9C27B0',
}

def persist_app_state(domain):
    """Save any unsaved LibreOffice state before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide5 = prs.slides[4]
    shapes = list(slide5.shapes)

    # Collect stage shapes (rounded rectangles with stage text), arrows, and description text boxes
    stage_shapes = {}  # label -> shape
    arrow_shapes = []
    desc_textboxes = []

    for shape in shapes:
        text = getattr(shape, 'text', '').strip()
        # Check if it's an auto shape
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # Try to get auto_shape_type for precise classification
            try:
                ast = shape.auto_shape_type
                # RIGHT_ARROW (33) or any arrow variant
                if ast == 33 or 'arrow' in shape.name.lower():
                    arrow_shapes.append(shape)
                    continue
            except (ValueError, AttributeError):
                pass
            # Check if it's a stage shape (rounded rectangle with a stage label)
            if text in EXPECTED_STAGES:
                stage_shapes[text] = shape
            # Also detect arrows by name if auto_shape_type check didn't catch it
            elif 'arrow' in shape.name.lower():
                arrow_shapes.append(shape)

        # Check for text boxes with multi-line description (not the title "Our Methodology")
        # Accept both TEXT_BOX and AUTO_SHAPE that have multi-line text and aren't stages/arrows
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and text and text != 'Our Methodology':
            lines = text.split('\n')
            if len(lines) >= 2:
                desc_textboxes.append(shape)

    # =========================================================================
    # Component 1: Four stage rounded rectangles with correct text (0.25 pts)
    # =========================================================================
    try:
        found_stages = [label for label in EXPECTED_STAGES if label in stage_shapes]
        stage_count = len(found_stages)
        if stage_count == 4:
            print(f"PASS: Component 1 — All 4 stages found: {found_stages} (0.25 pts)")
            total_score += 0.25
        elif stage_count > 0:
            partial = round(0.25 * (stage_count / 4), 2)
            print(f"PARTIAL: Component 1 — Found {stage_count}/4 stages: {found_stages} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — No stage shapes found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # =========================================================================
    # Component 2: Stage shapes have correct fill colors (0.20 pts)
    # =========================================================================
    try:
        color_matches = 0
        for label, expected_hex in EXPECTED_COLORS.items():
            if label not in stage_shapes:
                print(f"  SKIP: Component 2 — Stage '{label}' not found, cannot check color")
                continue
            shape = stage_shapes[label]
            try:
                fill = shape.fill
                if fill.type is not None and fill.type == 1:  # SOLID
                    actual_hex = str(fill.fore_color.rgb).upper()
                    expected_upper = expected_hex.upper()
                    if actual_hex == expected_upper:
                        color_matches += 1
                        print(f"  PASS: '{label}' fill color = {actual_hex}")
                    else:
                        print(f"  FAIL: '{label}' fill color = {actual_hex}, expected {expected_upper}")
                else:
                    print(f"  FAIL: '{label}' fill is not solid (type={fill.type})")
            except Exception as e:
                print(f"  ERROR: '{label}' fill check: {e}")

        if color_matches == 4:
            print(f"PASS: Component 2 — All 4 stage colors correct (0.20 pts)")
            total_score += 0.20
        elif color_matches > 0:
            partial = round(0.20 * (color_matches / 4), 2)
            print(f"PARTIAL: Component 2 — {color_matches}/4 colors correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No stage colors match")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # =========================================================================
    # Component 3: Three connecting arrows between stages (0.15 pts)
    # =========================================================================
    try:
        arrow_count = len(arrow_shapes)
        if arrow_count >= 3:
            print(f"PASS: Component 3 — Found {arrow_count} arrow shapes (0.15 pts)")
            total_score += 0.15
        elif arrow_count > 0:
            partial = round(0.15 * (arrow_count / 3), 2)
            print(f"PARTIAL: Component 3 — Found {arrow_count}/3 arrows ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No arrow shapes found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # =========================================================================
    # Component 4: Four 2-line description text boxes below stages (0.15 pts)
    # =========================================================================
    try:
        desc_count = len(desc_textboxes)
        if desc_count >= 4:
            for db in desc_textboxes[:4]:
                print(f"  DESC: '{db.text[:50]}...' at top={db.top}")
            print(f"PASS: Component 4 — Found {desc_count} description text boxes (0.15 pts)")
            total_score += 0.15
        elif desc_count > 0:
            partial = round(0.15 * (desc_count / 4), 2)
            print(f"PARTIAL: Component 4 — Found {desc_count}/4 descriptions ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — No description text boxes found")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # =========================================================================
    # Component 5: Stage shapes have Zoom animation (presetID=53) on click (0.15 pts)
    # =========================================================================
    try:
        # Parse animations from XML
        ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        animations = []  # list of (presetID, presetClass, nodeType, target_spid)

        with zipfile.ZipFile(file_path, 'r') as zf:
            try:
                with zf.open('ppt/slides/slide5.xml') as f:
                    root = ET.parse(f).getroot()
                for ctn in root.iter(f'{{{ns_p}}}cTn'):
                    preset_id = ctn.get('presetID')
                    preset_class = ctn.get('presetClass')
                    node_type = ctn.get('nodeType')
                    if preset_id:
                        for sp_tgt in ctn.iter(f'{{{ns_p}}}spTgt'):
                            spid = sp_tgt.get('spid')
                            animations.append((preset_id, preset_class, node_type, spid))
                            break
            except KeyError:
                print("  WARNING: Could not find slide5.xml in archive")

        # Build spid -> shape name mapping
        spid_map = {}
        for shape in shapes:
            spid_map[str(shape.shape_id)] = shape

        # Check that each stage shape has a Zoom (53) clickEffect animation
        stage_zoom_count = 0
        for label in EXPECTED_STAGES:
            if label not in stage_shapes:
                continue
            stage_spid = str(stage_shapes[label].shape_id)
            has_zoom = any(
                pid == '53' and pclass == 'entr' and ntype == 'clickEffect' and spid == stage_spid
                for pid, pclass, ntype, spid in animations
            )
            if has_zoom:
                stage_zoom_count += 1
                print(f"  PASS: '{label}' (spid={stage_spid}) has Zoom on click")
            else:
                print(f"  FAIL: '{label}' (spid={stage_spid}) missing Zoom on click animation")

        if stage_zoom_count == 4:
            print(f"PASS: Component 5 — All 4 stages have Zoom on click (0.15 pts)")
            total_score += 0.15
        elif stage_zoom_count > 0:
            partial = round(0.15 * (stage_zoom_count / 4), 2)
            print(f"PARTIAL: Component 5 — {stage_zoom_count}/4 stages with Zoom ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 — No stage shapes have Zoom animation")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # =========================================================================
    # Component 6: Arrows have Wipe (22) and descriptions have Fade (10) animations (0.10 pts)
    # =========================================================================
    try:
        # Check arrows have Wipe (presetID=22) animation
        arrow_spids = set(str(s.shape_id) for s in arrow_shapes)
        arrows_with_wipe = 0
        for spid in arrow_spids:
            has_wipe = any(
                pid == '22' and pclass == 'entr' and sp == spid
                for pid, pclass, ntype, sp in animations
            )
            if has_wipe:
                arrows_with_wipe += 1

        # Check descriptions have Fade (presetID=10) animation
        desc_spids = set(str(s.shape_id) for s in desc_textboxes)
        descs_with_fade = 0
        for spid in desc_spids:
            has_fade = any(
                pid == '10' and pclass == 'entr' and sp == spid
                for pid, pclass, ntype, sp in animations
            )
            if has_fade:
                descs_with_fade += 1

        # Score: half for arrows, half for descriptions
        arrow_sub = 0.0
        if len(arrow_spids) > 0 and arrows_with_wipe >= min(3, len(arrow_spids)):
            arrow_sub = 0.05
            print(f"  PASS: {arrows_with_wipe} arrows have Wipe animation")
        elif arrows_with_wipe > 0:
            arrow_sub = round(0.05 * (arrows_with_wipe / 3), 2)
            print(f"  PARTIAL: {arrows_with_wipe}/3 arrows have Wipe")
        else:
            print(f"  FAIL: No arrows have Wipe animation")

        desc_sub = 0.0
        if len(desc_spids) > 0 and descs_with_fade >= min(4, len(desc_spids)):
            desc_sub = 0.05
            print(f"  PASS: {descs_with_fade} descriptions have Fade animation")
        elif descs_with_fade > 0:
            desc_sub = round(0.05 * (descs_with_fade / 4), 2)
            print(f"  PARTIAL: {descs_with_fade}/4 descriptions have Fade")
        else:
            print(f"  FAIL: No descriptions have Fade animation")

        c6_score = arrow_sub + desc_sub
        if c6_score > 0:
            print(f"PASS: Component 6 — Arrow Wipe + Desc Fade ({c6_score} pts)")
            total_score += c6_score
        else:
            print(f"FAIL: Component 6 — No animation checks passed")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state('libreoffice_impress')

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
