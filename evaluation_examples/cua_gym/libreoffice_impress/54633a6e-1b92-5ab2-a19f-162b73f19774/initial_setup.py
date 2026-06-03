"""
Initial Setup: Create ProcessFlow_Pitch.pptx with 8 slides, slide 5 titled 'Our Methodology' with empty content.
Task ID: impress_sales_094
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_094'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_content(slide, title_text, body_lines):
    """Helper to populate a Title+Content slide."""
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line


def add_title_only(slide, title_text):
    """Set just the title on a Title Only layout slide."""
    slide.shapes.title.text = title_text


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "ProcessFlow Pitch"
    slide1.placeholders[1].text = "Transforming Ideas into Impact"

    # --- Slide 2: About Us ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_content(slide2, "About Meridian Consulting", [
        "Founded in 2018 by Sarah Chen and Marcus Rivera",
        "Headquartered in San Francisco with offices in London and Singapore",
        "Over 250 professionals serving Fortune 500 clients",
        "Specializing in digital transformation and process optimization",
        "Annual revenue exceeding $85M in fiscal year 2025",
    ])

    # --- Slide 3: The Challenge ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_content(slide3, "The Challenge", [
        "78% of enterprises struggle with legacy process inefficiencies",
        "Average time-to-market for new products: 14.3 months",
        "Customer satisfaction scores declining 12% year-over-year",
        "Digital adoption rates plateau at 45% without structured approach",
        "Competitor pressure increasing across all verticals",
    ])

    # --- Slide 4: Market Opportunity ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_content(slide4, "Market Opportunity", [
        "Global process consulting market valued at $42.3B in 2025",
        "Expected CAGR of 11.2% through 2030",
        "Enterprise SaaS integration demand growing 23% annually",
        "Mid-market segment underserved with only 18% penetration",
        "Regulatory compliance driving 35% of new engagements",
    ])

    # --- Slide 5: Our Methodology (EMPTY content - task target) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add just the title as a text box at the top
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Our Methodology"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    # Content area intentionally left empty for the task

    # --- Slide 6: Case Studies ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_content(slide6, "Case Studies", [
        "TechNova Corp: Reduced deployment cycles from 6 weeks to 8 days",
        "GlobalHealth Inc: Improved patient intake efficiency by 67%",
        "Pinnacle Financial: Automated 82% of compliance reporting",
        "EcoLogistics: Cut supply chain costs by $12.4M annually",
        "MediaWave: Increased content delivery speed by 340%",
    ])

    # --- Slide 7: Investment Ask ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_content(slide7, "Investment & Timeline", [
        "Phase 1 (Q1-Q2 2026): Discovery and baseline assessment - $180K",
        "Phase 2 (Q3 2026): Solution design and prototyping - $240K",
        "Phase 3 (Q4 2026): Pilot deployment across 3 business units - $320K",
        "Phase 4 (Q1 2027): Full rollout and optimization - $160K",
        "Total engagement: $900K with projected 4.2x ROI within 18 months",
    ])

    # --- Slide 8: Thank You ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[0])
    slide8.shapes.title.text = "Thank You"
    slide8.placeholders[1].text = "Questions? Contact us at partnerships@meridian.io"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
