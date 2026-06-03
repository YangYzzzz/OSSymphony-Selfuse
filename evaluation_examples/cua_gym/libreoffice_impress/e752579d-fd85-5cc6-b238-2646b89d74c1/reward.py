"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 67, the title’s underlined for some reason. I need to yank that underline and recolor the title text to Blue 4 (#000080) in LibreOffice Impress—what steps should I follow?
Generated: 2025-09-10 23:28:08
Status: success
Model: azure-o3
Total Steps: 2
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
import os

def verify_slide67_title_format(file_path: str) -> float:
    """Verify that on slide 67 the title text has:
    1. NO underline
    2. Font colour exactly Blue 4 (#000080)

    Returns a progressive score between 0.0 and 1.0.
    - 0.5 points for removing every underline
    - 0.5 points for colouring every run Blue 4
    """
    print(f"Verifying title formatting on slide 67 in: {file_path}")

    BLUE4 = RGBColor(0x00, 0x00, 0x80)  # Required colour
    total_score = 0.0

    # ----------------------------
    # 1  Load presentation file
    # ----------------------------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        return 0.0

    # -------------------------------------------------
    # 2  Ensure slide 67 exists (index 66, 0-based)
    # -------------------------------------------------
    if len(prs.slides) < 67:
        print(f"✗ Presentation only has {len(prs.slides)} slides – slide 67 missing")
        return 0.0
    slide = prs.slides[66]

    # -------------------------------------------------
    # 3  Locate the TITLE placeholder on that slide
    # -------------------------------------------------
    title_shape = None
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type == 1:  # 1 = TITLE placeholder
            title_shape = shape
            break
    if title_shape is None or not title_shape.has_text_frame:
        print("✗ No usable title placeholder found on slide 67")
        return 0.0

    # -------------------------------------------------
    # 4  Inspect every text run for underline & colour
    # -------------------------------------------------
    underline_ok = True
    colour_ok = True

    for para in title_shape.text_frame.paragraphs:
        for run in para.runs:
            # Underline check
            if run.font.underline not in (None, False):
                underline_ok = False
                print(f"✗ Underline still present in run '{run.text}' -> {run.font.underline}")
            # Colour check (if colour object exists)
            run_colour = run.font.color.rgb if run.font.color is not None else None
            if run_colour != BLUE4:
                colour_ok = False
                print(f"✗ Colour mismatch in run '{run.text}' -> {run_colour} (expected {BLUE4})")

    # ----------------------
    # 5  Progressive scoring
    # ----------------------
    if underline_ok:
        total_score += 0.5
        print("✓ All title runs have underline removed (0.5 points)")
    else:
        print("✗ Underline removal requirement NOT satisfied (0 points)")

    if colour_ok:
        total_score += 0.5
        print("✓ All title runs are Blue 4 (#000080) (0.5 points)")
    else:
        print("✗ Colour requirement NOT satisfied (0 points)")

    final_score = min(total_score, 1.0)
    print(f"Total score: {final_score}")
    return final_score


if __name__ == "__main__":
    FILE_PATH = "/home/user/on_slide_67_the_titles_underlined_for_some_reason_i_need_to_yank_that_underline_and_recolor_the_titl_golden.pptx"
    reward = verify_slide67_title_format(FILE_PATH)
    print(f"REWARD: {reward}")
