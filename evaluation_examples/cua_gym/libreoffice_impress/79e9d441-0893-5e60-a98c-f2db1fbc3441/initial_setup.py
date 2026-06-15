"""
Initial Setup: Create a 10-slide presentation with white backgrounds
Task ID: impress_el_086
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_el_086'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide content themes for a "Blue Theme" corporate presentation
    slide_data = [
        {
            "layout": 0,
            "title": "Blue Horizon Strategy 2025",
            "subtitle": "Annual Strategic Planning & Growth Initiative\nPresented by the Executive Leadership Team",
        },
        {
            "layout": 1,
            "title": "Market Overview",
            "body": "Global SaaS market projected to reach $908B by 2030\nCloud adoption accelerating across enterprise verticals\nAI-driven automation reshaping competitive landscape\nCustomer retention rates improving with personalized solutions",
        },
        {
            "layout": 1,
            "title": "Revenue Performance Q1-Q4",
            "body": "Q1: $12.4M (+18% YoY)\nQ2: $14.7M (+22% YoY)\nQ3: $15.9M (+25% YoY)\nQ4: $18.2M (+31% YoY)\nFull Year: $61.2M (+24% YoY)",
        },
        {
            "layout": 1,
            "title": "Product Roadmap",
            "body": "Phase 1: Core platform modernization (Complete)\nPhase 2: AI assistant integration (In Progress)\nPhase 3: Enterprise analytics dashboard (Q3 2025)\nPhase 4: Mobile-first experience redesign (Q4 2025)",
        },
        {
            "layout": 1,
            "title": "Customer Acquisition Funnel",
            "body": "Website Visitors: 2.4M monthly\nTrial Signups: 48,000 monthly\nActivated Users: 19,200 (40% activation)\nPaid Conversions: 5,760 (30% conversion)\nEnterprise Upgrades: 864 (15% upsell)",
        },
        {
            "layout": 1,
            "title": "Team Expansion Plan",
            "body": "Engineering: 45 current, hiring 20 more\nProduct Design: 12 current, hiring 6 more\nSales & Marketing: 30 current, hiring 15 more\nCustomer Success: 18 current, hiring 8 more\nTotal headcount target: 154 by end of 2025",
        },
        {
            "layout": 1,
            "title": "Competitive Analysis",
            "body": "Competitor A: Strong in enterprise, weak in SMB segment\nCompetitor B: Price leader but limited feature set\nCompetitor C: Best UX but poor scalability\nOur advantage: Full-stack solution with AI differentiation",
        },
        {
            "layout": 1,
            "title": "Financial Projections",
            "body": "2025 Revenue Target: $82M (+34% growth)\nGross Margin: 78% (up from 72%)\nOperating Expenses: $54M\nEBITDA Target: $10M (first profitable year)\nCash Runway: 36+ months",
        },
        {
            "layout": 1,
            "title": "Risk Assessment",
            "body": "Market Risk: Economic downturn affecting enterprise budgets\nTechnology Risk: Rapid AI advancement requiring pivot\nTalent Risk: Competitive hiring in engineering\nRegulatory Risk: Data privacy compliance across regions\nMitigation strategies documented in appendix",
        },
        {
            "layout": 1,
            "title": "Next Steps & Action Items",
            "body": "Complete Series C fundraising by June 2025\nLaunch AI assistant beta to top 50 customers\nExpand EMEA sales team by Q3\nAchieve SOC 2 Type II certification\nHost annual customer summit in September",
        },
    ]

    for i, sd in enumerate(slide_data):
        layout_idx = sd["layout"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Set explicit white background on every slide
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = sd["title"]
            for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                run.font.size = Pt(36)
                run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        # Set body / subtitle
        body_key = "subtitle" if layout_idx == 0 else "body"
        if body_key in sd and len(slide.placeholders) > 1:
            ph = slide.placeholders[1]
            ph.text = ""
            lines = sd[body_key].split("\n")
            for j, line in enumerate(lines):
                if j == 0:
                    ph.text_frame.paragraphs[0].text = line
                    para = ph.text_frame.paragraphs[0]
                else:
                    para = ph.text_frame.add_paragraph()
                    para.text = line
                para.space_after = Pt(6)
                for run in para.runs:
                    run.font.size = Pt(18) if layout_idx != 0 else Pt(20)
                    run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
