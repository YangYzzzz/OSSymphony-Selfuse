"""
Reward Script: Digital Marketing Analytics Presentation
Task ID: impress_wf_079
Domain: libreoffice_impress
Scoring:
  C1  (0.15) - File exists as .pptx on Desktop with exactly 12 slides
  C2  (0.10) - Slide 1 title text contains 'Digital Marketing Performance Q3'
  C3  (0.10) - Slide 2 has a doughnut/pie chart (channel mix)
  C4  (0.10) - Slide 3 has metric cards (sessions/bounce/duration) + line chart
  C5  (0.10) - Slide 4 has SEO keyword rankings table
  C6  (0.10) - Slide 5 has PPC table + chart
  C7  (0.05) - Slide 6 has social media chart
  C8  (0.10) - Slide 7 has funnel shapes with decreasing widths
  C9  (0.05) - Slide 8 has content performance table
  C10 (0.05) - Slide 9 has pie chart (attribution model)
  C11 (0.05) - Slides 10-11 have charts (CAC trend + ROI by channel)
  C12 (0.05) - Theme colors #6A1B9A and #00897B used
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_079'
FILE_PATH = os.path.join(WORKDIR, 'Desktop', 'Marketing_Analytics.pptx')


def get_all_text(slide):
    """Collect all text from a slide's shapes (including nested groups)."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    return texts


def has_chart_of_type(slide, chart_types):
    """Check if slide contains a chart matching any of the given type codes."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            try:
                if shape.chart.chart_type in chart_types:
                    return True
            except Exception:
                pass
    return False


def has_any_chart(slide):
    """Check if slide has any chart."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            return True
    return False


def has_table(slide, min_rows=2, min_cols=2):
    """Check if slide has a table with minimum dimensions."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            t = shape.table
            if len(t.rows) >= min_rows and len(t.columns) >= min_cols:
                return True
    return False


def count_auto_shapes_with_text(slide):
    """Count auto shapes that contain text (metric cards, funnel steps)."""
    count = 0
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.has_text_frame:
            if shape.text_frame.text.strip():
                count += 1
    return count


def check_theme_colors(prs):
    """Check if #6A1B9A (purple) and #00897B (teal) are used anywhere."""
    found_purple = False
    found_teal = False
    for slide in prs.slides:
        for shape in slide.shapes:
            # Check text colors
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color.type is not None:
                                rgb = str(run.font.color.rgb)
                                if rgb == '6A1B9A':
                                    found_purple = True
                                if rgb == '00897B':
                                    found_teal = True
                        except Exception:
                            pass
            # Check shape fill colors
            try:
                fill = shape.fill
                if fill.type == 1:  # solid fill
                    rgb = str(fill.fore_color.rgb)
                    if rgb == '6A1B9A':
                        found_purple = True
                    if rgb == '00897B':
                        found_teal = True
            except Exception:
                pass
        if found_purple and found_teal:
            break
    return found_purple, found_teal


def verify_task():
    """Verify task completion with progressive scoring. Returns float 0.0-1.0."""
    total_score = 0.0

    # Precondition: file must exist
    if not os.path.exists(FILE_PATH):
        print(f"CRITICAL: File not found: {FILE_PATH}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(FILE_PATH)
    except Exception as e:
        print(f"CRITICAL: Cannot load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)

    # Component 1: Exactly 12 slides (0.15 pts)
    try:
        slide_count = len(slides)
        if slide_count == 12:
            print(f"PASS: C1 - Presentation has exactly 12 slides (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: C1 - Expected 12 slides, found {slide_count}")
    except Exception as e:
        print(f"ERROR: C1 - {e}")

    # Guard: need at least 12 slides for remaining checks
    if len(slides) < 12:
        print(f"WARNING: Only {len(slides)} slides, some checks will be skipped")
        final_score = min(total_score, 1.0)
        print(f"\nScore: {final_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: Slide 1 title "Digital Marketing Performance Q3" (0.10 pts)
    try:
        slide1_texts = get_all_text(slides[0])
        combined = ' '.join(slide1_texts).lower()
        if 'digital marketing performance q3' in combined:
            print(f"PASS: C2 - Slide 1 has title 'Digital Marketing Performance Q3' (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: C2 - Slide 1 title not found. Texts: {slide1_texts[:3]}")
    except Exception as e:
        print(f"ERROR: C2 - {e}")

    # Component 3: Slide 2 has doughnut chart (channel mix) (0.10 pts)
    # Doughnut chart type is -4120, also accept PIE (5) as donut variant
    try:
        has_donut = False
        for shape in slides[1].shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                ct = shape.chart.chart_type
                # DOUGHNUT = -4120, DOUGHNUT_EXPLODED = -4121, PIE = 5
                if ct in (-4120, -4121, 5):
                    has_donut = True
                    break
        if has_donut:
            print(f"PASS: C3 - Slide 2 has doughnut/pie chart for channel mix (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: C3 - Slide 2 has no doughnut/pie chart")
    except Exception as e:
        print(f"ERROR: C3 - {e}")

    # Component 4: Slide 3 has metric cards (auto shapes with text) + line chart (0.10 pts)
    try:
        card_count = count_auto_shapes_with_text(slides[2])
        has_line = False
        for shape in slides[2].shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                ct = shape.chart.chart_type
                # LINE = 4, LINE_MARKERS = 65, LINE_STACKED = 63, etc.
                if ct in (4, 63, 64, 65, 66, 67, 68):
                    has_line = True
                    break
        if card_count >= 3 and has_line:
            print(f"PASS: C4 - Slide 3 has {card_count} metric cards + line chart (0.10 pts)")
            total_score += 0.10
        elif card_count >= 3:
            print(f"PARTIAL: C4 - Slide 3 has {card_count} metric cards but no line chart (0.05 pts)")
            total_score += 0.05
        elif has_line:
            print(f"PARTIAL: C4 - Slide 3 has line chart but only {card_count} metric cards (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: C4 - Slide 3: {card_count} metric cards, line chart={has_line}")
    except Exception as e:
        print(f"ERROR: C4 - {e}")

    # Component 5: Slide 4 has SEO keyword rankings table (0.10 pts)
    try:
        if has_table(slides[3], min_rows=3, min_cols=3):
            print(f"PASS: C5 - Slide 4 has SEO keyword rankings table (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: C5 - Slide 4 has no suitable table")
    except Exception as e:
        print(f"ERROR: C5 - {e}")

    # Component 6: Slide 5 has PPC table + chart (0.10 pts)
    try:
        ppc_table = has_table(slides[4], min_rows=3, min_cols=4)
        ppc_chart = has_any_chart(slides[4])
        if ppc_table and ppc_chart:
            print(f"PASS: C6 - Slide 5 has PPC table + chart (0.10 pts)")
            total_score += 0.10
        elif ppc_table or ppc_chart:
            print(f"PARTIAL: C6 - Slide 5 has table={ppc_table}, chart={ppc_chart} (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: C6 - Slide 5 has no PPC table or chart")
    except Exception as e:
        print(f"ERROR: C6 - {e}")

    # Component 7: Slide 6 has social media chart (0.05 pts)
    try:
        if has_any_chart(slides[5]):
            print(f"PASS: C7 - Slide 6 has social media engagement chart (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: C7 - Slide 6 has no chart")
    except Exception as e:
        print(f"ERROR: C7 - {e}")

    # Component 8: Slide 7 has funnel shapes with decreasing widths (0.10 pts)
    try:
        funnel_shapes = []
        for shape in slides[6].shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    funnel_shapes.append((shape.width, txt))
        # Need at least 3 funnel shapes and first should be wider than last
        if len(funnel_shapes) >= 3:
            widths = [w for w, _ in funnel_shapes]
            # Check that widths generally decrease (first > last)
            if widths[0] > widths[-1]:
                print(f"PASS: C8 - Slide 7 has {len(funnel_shapes)} funnel shapes, widths decrease ({widths[0]} > {widths[-1]}) (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: C8 - Funnel shapes don't decrease in width: {widths}")
        else:
            print(f"FAIL: C8 - Only {len(funnel_shapes)} auto shapes with text on slide 7, need >=3")
    except Exception as e:
        print(f"ERROR: C8 - {e}")

    # Component 9: Slide 8 has content performance table (0.05 pts)
    try:
        if has_table(slides[7], min_rows=3, min_cols=3):
            print(f"PASS: C9 - Slide 8 has content performance table (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: C9 - Slide 8 has no suitable table")
    except Exception as e:
        print(f"ERROR: C9 - {e}")

    # Component 10: Slide 9 has pie chart (attribution model) (0.05 pts)
    try:
        has_pie = False
        for shape in slides[8].shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                ct = shape.chart.chart_type
                # PIE = 5, PIE_EXPLODED = -4102, DOUGHNUT = -4120
                if ct in (5, -4102, -4120, -4121):
                    has_pie = True
                    break
        if has_pie:
            print(f"PASS: C10 - Slide 9 has pie/doughnut chart for attribution (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: C10 - Slide 9 has no pie chart")
    except Exception as e:
        print(f"ERROR: C10 - {e}")

    # Component 11: Slides 10-11 have charts (CAC trend + ROI by channel) (0.05 pts)
    try:
        s10_chart = has_any_chart(slides[9])
        s11_chart = has_any_chart(slides[10])
        if s10_chart and s11_chart:
            print(f"PASS: C11 - Slides 10-11 both have charts (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: C11 - Slide 10 chart={s10_chart}, Slide 11 chart={s11_chart}")
    except Exception as e:
        print(f"ERROR: C11 - {e}")

    # Component 12: Theme colors #6A1B9A and #00897B used (0.05 pts)
    try:
        found_purple, found_teal = check_theme_colors(prs)
        if found_purple and found_teal:
            print(f"PASS: C12 - Both theme colors #6A1B9A and #00897B found (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: C12 - Purple(#6A1B9A)={found_purple}, Teal(#00897B)={found_teal}")
    except Exception as e:
        print(f"ERROR: C12 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


verify_task()
