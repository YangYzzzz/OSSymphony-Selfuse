"""
Reward Script: Merge first row of table on slide 4 into spanning header
Task ID: impress_tct_005
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): First cell gridSpan == 4 (merged across all columns)
  Component 2 (0.3): Merged cell text is 'Annual Performance Summary'
  Component 3 (0.2): Continuation cells (0,1)-(0,3) have hMerge attribute
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_005'


def persist_app_state(domain: str):
    """Save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 4 slides
    if len(prs.slides) < 4:
        print(f"CRITICAL: Expected at least 4 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)

    # Find the table on slide 4
    table = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            break

    if table is None:
        print("CRITICAL: No table found on slide 4")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: table has expected dimensions
    if len(table.rows) < 1 or len(table.columns) < 4:
        print(f"CRITICAL: Table dimensions unexpected: {len(table.rows)} rows x {len(table.columns)} cols")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: First cell (0,0) has gridSpan == 4 (0.5 points)
    # This verifies the merge spans all 4 columns
    try:
        tc_0_0 = table.cell(0, 0)._tc
        grid_span = int(tc_0_0.get('gridSpan', '1'))
        num_cols = len(table.columns)
        if grid_span >= num_cols:
            print(f"PASS: Component 1 -- Cell(0,0) gridSpan={grid_span} spans all {num_cols} columns (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 1 -- Cell(0,0) gridSpan={grid_span}, expected {num_cols}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Merged cell text is 'Annual Performance Summary' (0.3 points)
    # This verifies the merged header contains the correct text
    try:
        cell_text = table.cell(0, 0).text.strip()
        expected_text = 'Annual Performance Summary'
        if cell_text == expected_text:
            print(f"PASS: Component 2 -- Cell(0,0) text is '{cell_text}' (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 -- Cell(0,0) text is '{cell_text}', expected '{expected_text}'")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Continuation cells (0,1)-(0,3) have hMerge attribute (0.2 points)
    # In a merged row, non-origin cells should have hMerge='1' or hMerge=True
    try:
        hmerge_count = 0
        for c in range(1, len(table.columns)):
            tc = table.cell(0, c)._tc
            hmerge = tc.get('hMerge')
            if hmerge is not None:
                hmerge_count += 1

        expected_hmerge = len(table.columns) - 1  # columns 1,2,3
        if hmerge_count == expected_hmerge:
            print(f"PASS: Component 3 -- {hmerge_count}/{expected_hmerge} continuation cells have hMerge (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 3 -- Only {hmerge_count}/{expected_hmerge} continuation cells have hMerge")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = round(min(total_score, 1.0), 1)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist app state before verification
persist_app_state("libreoffice_impress")

# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
