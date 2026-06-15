"""
Initial Setup: Performance Review presentation with a table on slide 4
Task ID: impress_tct_005
Domain: libreoffice_impress

Creates a 6-slide presentation. Slide 4 has a 4-column x 6-row table where the
first row has four SEPARATE cells: "Annual", "Performance", "Summary", and empty.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_005'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_cell_text(cell, text, font_name="Calibri", font_size=Pt(11),
                  bold=False, color=None, alignment=None):
    """Helper to set text and formatting on a table cell."""
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if alignment:
        p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Annual Performance Review 2025"
    slide1.placeholders[1].text = "Prepared by Human Resources Department"

    # ---- Slide 2: Overview ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Review Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "This presentation covers the annual performance metrics for all departments."
    p2 = body2.add_paragraph()
    p2.text = "Key areas evaluated include productivity, collaboration, innovation, and leadership."
    p3 = body2.add_paragraph()
    p3.text = "Data collected from Q1-Q4 2025 across 12 regional offices."
    p4 = body2.add_paragraph()
    p4.text = "Total employees reviewed: 847"

    # ---- Slide 3: Key Metrics ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Key Performance Indicators"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Average employee satisfaction score: 4.2/5.0"
    for line in [
        "Revenue per employee: $142,300",
        "Training hours completed: 32,450",
        "Project completion rate: 91.7%",
        "Employee retention rate: 88.3%",
    ]:
        p = body3.add_paragraph()
        p.text = line

    # ---- Slide 4: Performance Table (the target slide) ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title text box
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Department Performance Data"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Add 4-column x 6-row table
    rows, cols = 6, 4
    table_shape = slide4.shapes.add_table(
        rows, cols,
        Inches(0.8), Inches(1.4),
        Inches(10.5), Inches(4.5)
    )
    table = table_shape.table

    # Set column widths
    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.5)
    table.columns[3].width = Inches(2.7)

    # Row 0: Header row - four SEPARATE cells (NOT merged)
    header_color = RGBColor(0xFF, 0xFF, 0xFF)
    header_bg = RGBColor(0x1F, 0x49, 0x7D)

    header_texts = ["Annual", "Performance", "Summary", ""]
    for c, txt in enumerate(header_texts):
        cell = table.cell(0, c)
        set_cell_text(cell, txt, font_size=Pt(14), bold=True,
                      color=header_color, alignment=PP_ALIGN.CENTER)
        # Set background
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = header_bg

    # Row 1: Sub-headers
    sub_headers = ["Department", "Score (Q1-Q4)", "Rating", "Trend"]
    sub_bg = RGBColor(0x4A, 0x7F, 0xB5)
    for c, txt in enumerate(sub_headers):
        cell = table.cell(1, c)
        set_cell_text(cell, txt, font_size=Pt(12), bold=True,
                      color=header_color, alignment=PP_ALIGN.CENTER)
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = sub_bg

    # Rows 2-5: Data rows
    data_rows = [
        ["Engineering", "92.4", "Exceeds Expectations", "Improving"],
        ["Marketing", "87.1", "Meets Expectations", "Stable"],
        ["Sales", "94.8", "Exceeds Expectations", "Improving"],
        ["Operations", "83.6", "Meets Expectations", "Needs Attention"],
    ]
    alt_colors = [RGBColor(0xE8, 0xF0, 0xFE), RGBColor(0xFF, 0xFF, 0xFF)]
    for r_idx, row_data in enumerate(data_rows):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 2, c_idx)
            set_cell_text(cell, val, font_size=Pt(11),
                          alignment=PP_ALIGN.CENTER)
            cell_fill = cell.fill
            cell_fill.solid()
            cell_fill.fore_color.rgb = alt_colors[r_idx % 2]

    # ---- Slide 5: Goals ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Goals for 2026"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Increase overall employee satisfaction to 4.5/5.0"
    for line in [
        "Reduce voluntary turnover by 5%",
        "Launch new leadership development program by Q2",
        "Achieve 95% project completion rate across all departments",
        "Expand mentorship program to cover all junior employees",
    ]:
        p = body5.add_paragraph()
        p.text = line

    # ---- Slide 6: Closing ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[0])
    slide6.shapes.title.text = "Thank You"
    slide6.placeholders[1].text = "Questions? Contact HR at hr@company.com"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
