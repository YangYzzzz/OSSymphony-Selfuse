"""
FINAL REWARD SCRIPT - SUCCESS
Task: Every time I tweak something, I end up copying the header to each slide one by one – super tedious. How do I set the text "Q3 Review" so it shows up automatically at the top of every slide in my LibreOffice Impress deck?
Generated: 2025-09-10 15:07:12
Status: success
Model: azure-o3
Total Steps: 1
"""

from pptx import Presentation
import os


def verify_header_on_all_slides(file_path: str, header_text: str = "Q3 Review") -> float:
    """Verify that the given presentation shows the required header text on every slide.

    Scoring (progressive):
    1. Up to 0.8 points if the text appears on every slide (pro-rated if not).
    2. Up to 0.2 additional points if that text is positioned near the top of
       every slide (≤ 2 inches from the top, pro-rated if not).

    Returns a float between 0.0 and 1.0 and prints detailed diagnostics.
    """

    max_score = 1.0
    score = 0.0

    # ---- 1. File existence & loading (no points, prerequisite) ----
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Could not open presentation: {e}")
        return 0.0

    if not prs.slides:
        print("✗ Presentation contains no slides")
        return 0.0

    total_slides = len(prs.slides)
    header_text_lower = header_text.lower()

    # ---- 2. Analyse every slide ----
    slides_with_header = 0           # header text present anywhere
    slides_with_header_at_top = 0    # header text present near top (≤ 2")

    for idx, slide in enumerate(prs.slides, start=1):
        found_anywhere = False
        found_up_top = False

        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue  # skip non-text shapes

            text = shape.text.strip()
            if not text:
                continue

            if header_text_lower in text.lower():
                found_anywhere = True
                # location check if shape has coordinate attributes
                if hasattr(shape, "top"):
                    try:
                        # ≤ 2 inches (2 * 914400 EMUs) from top of slide
                        if shape.top <= 2 * 914400:
                            found_up_top = True
                    except Exception:
                        pass  # ignore shapes without valid coordinates

        # Update counters for this slide
        if found_anywhere:
            slides_with_header += 1
        if found_up_top:
            slides_with_header_at_top += 1

    # ---- 3. Progressive scoring ----
    header_ratio = slides_with_header / total_slides      # between 0 and 1
    top_ratio = slides_with_header_at_top / total_slides  # between 0 and 1

    score += header_ratio * 0.8  # up to 0.8 based on presence across slides
    score += top_ratio * 0.2     # up to 0.2 based on correct positioning

    # ---- 4. Diagnostics ----
    print(f"Slides with header text  : {slides_with_header}/{total_slides} (ratio {header_ratio:.2f})")
    print(f"Slides with header @ top : {slides_with_header_at_top}/{total_slides} (ratio {top_ratio:.2f})")
    print(f"Partial scores -> presence: {header_ratio * 0.8:.2f}, position: {top_ratio * 0.2:.2f}")

    final_score = min(score, max_score)
    print(f"Final score: {final_score}")
    return final_score


if __name__ == "__main__":
    # Path provided by the task context
    file_path = "/home/user/every_time_i_tweak_something_i_end_up_copying_the_header_to_each_slide_one_by_one_super_tedious_how__golden.pptx"

    reward = verify_header_on_all_slides(file_path)
    print(f"REWARD: {reward}")

