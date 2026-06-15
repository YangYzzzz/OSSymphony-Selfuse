"""
Initial Setup: Product catalog presentation with image on slide 2
Task ID: impress_tm_052
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_052'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
IMG_PATH = f'{WORKDIR}/product_image.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_product_image():
    """Create a simple product image for the catalog."""
    img = Image.new('RGB', (600, 400), color=(41, 128, 185))
    # Draw some simple rectangles to simulate a product photo
    from PIL import ImageDraw, ImageFont
    draw = ImageDraw.Draw(img)
    # Product box
    draw.rectangle([100, 50, 500, 350], fill=(52, 152, 219), outline=(255, 255, 255), width=3)
    # Label area
    draw.rectangle([150, 150, 450, 250], fill=(236, 240, 241))
    # Text
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 28)
        draw.text((180, 175), "SmartWidget", fill=(44, 62, 80), font=font)
    except (IOError, OSError):
        draw.text((200, 185), "SmartWidget", fill=(44, 62, 80))
    img.save(IMG_PATH)
    print(f'Product image created: {IMG_PATH}')


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    # Standard 16:9 slide dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "TechNova Product Catalog"
    slide1.placeholders[1].text = "Q2 2025 Edition"

    # --- Slide 2: Product Spotlight (with centered image 15cm x 10cm) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_text_box(slide2, Cm(1), Cm(0.5), Cm(20), Cm(2),
                 "Product Spotlight: SmartWidget Pro", font_size=28, bold=True,
                 color=RGBColor(0x2C, 0x3E, 0x50))

    # Image: 15cm x 10cm, centered on slide
    img_w = Cm(15)
    img_h = Cm(10)
    img_left = (slide_w - img_w) // 2
    img_top = (slide_h - img_h) // 2
    slide2.shapes.add_picture(IMG_PATH, img_left, img_top, img_w, img_h)

    # --- Slide 3: Product Features ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide3, Cm(1), Cm(0.5), Cm(25), Cm(2),
                 "Key Features", font_size=28, bold=True,
                 color=RGBColor(0x2C, 0x3E, 0x50))
    features = [
        "Wireless connectivity with Bluetooth 5.3 and WiFi 6E",
        "12-hour battery life with fast charging support",
        "IP68 water and dust resistance rating",
        "AI-powered noise cancellation technology",
        "Compact design at only 45g weight"
    ]
    y_pos = Cm(3)
    for feat in features:
        add_text_box(slide3, Cm(2), y_pos, Cm(28), Cm(1.2),
                     f"• {feat}", font_size=16)
        y_pos += Cm(1.8)

    # --- Slide 4: Pricing Table ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide4, Cm(1), Cm(0.5), Cm(20), Cm(2),
                 "Pricing & Plans", font_size=28, bold=True,
                 color=RGBColor(0x2C, 0x3E, 0x50))
    table_shape = slide4.shapes.add_table(4, 3, Cm(3), Cm(3.5), Cm(26), Cm(8))
    table = table_shape.table
    headers = ["Plan", "Price", "Features"]
    data = [
        ["Basic", "$29.99/mo", "Core features, 1 device"],
        ["Professional", "$59.99/mo", "All features, 5 devices"],
        ["Enterprise", "$149.99/mo", "Unlimited devices, priority support"]
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 5: Customer Testimonials ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide5, Cm(1), Cm(0.5), Cm(25), Cm(2),
                 "What Our Customers Say", font_size=28, bold=True,
                 color=RGBColor(0x2C, 0x3E, 0x50))
    testimonials = [
        ('"The SmartWidget Pro transformed our workflow. Best investment this year."',
         "— Sarah Chen, VP of Operations, Meridian Corp"),
        ('"Incredible battery life and seamless integration with our existing tools."',
         "— Marcus Johnson, CTO, Brightpath Solutions"),
    ]
    y_pos = Cm(3.5)
    for quote, author in testimonials:
        add_text_box(slide5, Cm(2), y_pos, Cm(28), Cm(2),
                     quote, font_size=16, color=RGBColor(0x34, 0x49, 0x5E))
        add_text_box(slide5, Cm(3), y_pos + Cm(2), Cm(25), Cm(1),
                     author, font_size=13, bold=True,
                     color=RGBColor(0x7F, 0x8C, 0x8D))
        y_pos += Cm(4)

    # --- Slide 6: Contact / CTA ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide6, Cm(3), Cm(2), Cm(28), Cm(2),
                 "Ready to Get Started?", font_size=36, bold=True,
                 color=RGBColor(0x2C, 0x3E, 0x50), alignment=PP_ALIGN.CENTER)
    add_text_box(slide6, Cm(3), Cm(5), Cm(28), Cm(1.5),
                 "Contact our sales team at sales@technova.io", font_size=20,
                 color=RGBColor(0x7F, 0x8C, 0x8D), alignment=PP_ALIGN.CENTER)
    add_text_box(slide6, Cm(3), Cm(7), Cm(28), Cm(1.5),
                 "Visit us: www.technova.io/smartwidget", font_size=18,
                 color=RGBColor(0x29, 0x80, 0xB9), alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_product_image()
create_initial()
