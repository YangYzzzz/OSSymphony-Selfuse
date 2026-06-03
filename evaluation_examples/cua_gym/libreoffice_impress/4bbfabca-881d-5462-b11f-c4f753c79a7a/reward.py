"""
Reward Script: Resize image on slide 2 to 10cm x 7cm and position at top-right corner
Task ID: impress_tm_052
Domain: libreoffice_impress
Scoring:
  Component 1: Image width is 10cm (0.3 pts)
  Component 2: Image height is 7cm (0.3 pts)
  Component 3: Image positioned at top-right corner of slide (0.4 pts)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_052'

# Conversion constants
CM_TO_EMU = 360000  # 1 cm = 360000 EMU

# Expected values
EXPECTED_WIDTH_CM = 10.0
EXPECTED_HEIGHT_CM = 7.0
EXPECTED_WIDTH_EMU = int(EXPECTED_WIDTH_CM * CM_TO_EMU)   # 3600000
EXPECTED_HEIGHT_EMU = int(EXPECTED_HEIGHT_CM * CM_TO_EMU)  # 2520000

# Tolerance: 2% relative for size, absolute 0.5cm for position
SIZE_TOLERANCE = 0.02
POSITION_TOLERANCE_EMU = int(0.5 * CM_TO_EMU)  # 180000 EMU = 0.5cm


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 2")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]
    slide_width = prs.slide_width

    # Find the picture shape on slide 2
    picture = None
    for shape in slide2.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture = shape
            break

    if picture is None:
        print("FAIL: No picture found on slide 2")
        print("REWARD: 0.0")
        return 0.0

    actual_width = picture.width
    actual_height = picture.height
    actual_left = picture.left
    actual_top = picture.top

    print(f"Image properties: width={actual_width} EMU ({actual_width/CM_TO_EMU:.2f}cm), "
          f"height={actual_height} EMU ({actual_height/CM_TO_EMU:.2f}cm)")
    print(f"Image position: left={actual_left} EMU ({actual_left/CM_TO_EMU:.2f}cm), "
          f"top={actual_top} EMU ({actual_top/CM_TO_EMU:.2f}cm)")
    print(f"Slide width: {slide_width} EMU ({slide_width/CM_TO_EMU:.2f}cm)")

    # Component 1: Image width is 10cm (0.3 points)
    try:
        width_diff = abs(actual_width - EXPECTED_WIDTH_EMU)
        width_tolerance = EXPECTED_WIDTH_EMU * SIZE_TOLERANCE
        if width_diff <= width_tolerance:
            print(f"PASS: Component 1 -- Image width is {actual_width/CM_TO_EMU:.2f}cm "
                  f"(expected {EXPECTED_WIDTH_CM}cm, diff={width_diff/CM_TO_EMU:.3f}cm) (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 -- Image width is {actual_width/CM_TO_EMU:.2f}cm, "
                  f"expected {EXPECTED_WIDTH_CM}cm (diff={width_diff/CM_TO_EMU:.3f}cm)")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Image height is 7cm (0.3 points)
    try:
        height_diff = abs(actual_height - EXPECTED_HEIGHT_EMU)
        height_tolerance = EXPECTED_HEIGHT_EMU * SIZE_TOLERANCE
        if height_diff <= height_tolerance:
            print(f"PASS: Component 2 -- Image height is {actual_height/CM_TO_EMU:.2f}cm "
                  f"(expected {EXPECTED_HEIGHT_CM}cm, diff={height_diff/CM_TO_EMU:.3f}cm) (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 -- Image height is {actual_height/CM_TO_EMU:.2f}cm, "
                  f"expected {EXPECTED_HEIGHT_CM}cm (diff={height_diff/CM_TO_EMU:.3f}cm)")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Image positioned at top-right corner (0.4 points)
    # Top-right means: right edge aligned with slide right edge, top edge near slide top
    try:
        right_edge = actual_left + actual_width
        right_gap = abs(slide_width - right_edge)  # distance from right edge to slide right

        # Top edge should be near 0 (top of slide)
        top_gap = actual_top  # distance from top

        top_ok = top_gap <= POSITION_TOLERANCE_EMU
        right_ok = right_gap <= POSITION_TOLERANCE_EMU

        print(f"  Right edge gap: {right_gap/CM_TO_EMU:.2f}cm (tolerance: {POSITION_TOLERANCE_EMU/CM_TO_EMU:.2f}cm)")
        print(f"  Top edge gap: {top_gap/CM_TO_EMU:.2f}cm (tolerance: {POSITION_TOLERANCE_EMU/CM_TO_EMU:.2f}cm)")

        if top_ok and right_ok:
            print(f"PASS: Component 3 -- Image is at top-right corner "
                  f"(right gap={right_gap/CM_TO_EMU:.2f}cm, top gap={top_gap/CM_TO_EMU:.2f}cm) (0.4 pts)")
            total_score += 0.4
        else:
            reasons = []
            if not right_ok:
                reasons.append(f"right gap={right_gap/CM_TO_EMU:.2f}cm")
            if not top_ok:
                reasons.append(f"top gap={top_gap/CM_TO_EMU:.2f}cm")
            print(f"FAIL: Component 3 -- Image NOT at top-right corner: {', '.join(reasons)}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
