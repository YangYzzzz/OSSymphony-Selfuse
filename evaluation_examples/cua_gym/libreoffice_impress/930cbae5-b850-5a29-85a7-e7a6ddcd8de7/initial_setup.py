"""
Initial Setup: 5-slide internal strategy deck — plain titles on slides 2 and 3,
no CONFIDENTIAL text on slide 1.
Task ID: osworld_impress_multi_op_combined_005
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_multi_op_combined_005'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard widescreen slide dimensions (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_layouts = prs.slide_layouts
    # Layout 0 = Title Slide, Layout 1 = Title and Content, Layout 5 = Blank

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(slide_layouts[0])
    slide1.shapes.title.text = "Q3 2025 Strategic Planning"
    # subtitle placeholder
    for ph in slide1.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = "Internal Use Only — Executive Team"
            break

    # ---- Slide 2: Market Analysis (plain title — no underline) ----
    slide2 = prs.slides.add_slide(slide_layouts[1])
    slide2.shapes.title.text = "Market Analysis Overview"
    for ph in slide2.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.text = "Current market landscape"
            p2 = tf.add_paragraph()
            p2.text = "Competitor benchmarking across 12 key verticals"
            p3 = tf.add_paragraph()
            p3.text = "Addressable market estimated at $4.2B by 2027"
            p4 = tf.add_paragraph()
            p4.text = "Customer acquisition cost decreased 18% YoY"
            break

    # ---- Slide 3: Growth Strategy (plain title — no underline) ----
    slide3 = prs.slides.add_slide(slide_layouts[1])
    slide3.shapes.title.text = "Growth Strategy Roadmap"
    for ph in slide3.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.text = "Phase 1: Core product expansion (Q3–Q4 2025)"
            p2 = tf.add_paragraph()
            p2.text = "Phase 2: International market entry (Q1 2026)"
            p3 = tf.add_paragraph()
            p3.text = "Phase 3: Strategic acquisitions pipeline"
            p4 = tf.add_paragraph()
            p4.text = "Target: 35% revenue growth by end of FY2026"
            break

    # ---- Slide 4: Financial Projections ----
    slide4 = prs.slides.add_slide(slide_layouts[1])
    slide4.shapes.title.text = "Financial Projections FY2026"
    for ph in slide4.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.text = "Revenue forecast: $128M (+35% YoY)"
            p2 = tf.add_paragraph()
            p2.text = "EBITDA margin target: 22%"
            p3 = tf.add_paragraph()
            p3.text = "CapEx budget: $18M for infrastructure expansion"
            p4 = tf.add_paragraph()
            p4.text = "R&D investment: $9.5M (7.4% of revenue)"
            break

    # ---- Slide 5: Next Steps ----
    slide5 = prs.slides.add_slide(slide_layouts[1])
    slide5.shapes.title.text = "Action Items & Next Steps"
    for ph in slide5.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.text = "Finalize Q3 resource allocation by July 31"
            p2 = tf.add_paragraph()
            p2.text = "Conduct market entry feasibility study — APAC region"
            p3 = tf.add_paragraph()
            p3.text = "Review partnership proposals from 3 strategic vendors"
            p4 = tf.add_paragraph()
            p4.text = "Board presentation scheduled for August 12, 2025"
            break

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
