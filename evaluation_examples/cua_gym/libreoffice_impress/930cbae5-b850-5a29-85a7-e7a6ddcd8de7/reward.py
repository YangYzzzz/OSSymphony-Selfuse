"""
Reward Script: Apply underline formatting to titles on slides 2 and 3,
               and add 'CONFIDENTIAL' text element at the bottom of slide 1.
Task ID: osworld_impress_multi_op_combined_005
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.35): Slide 2 title has underline formatting
  - Component 2 (0.35): Slide 3 title has underline formatting
  - Component 3 (0.30): Slide 1 has a text element containing 'CONFIDENTIAL' at the bottom area
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_multi_op_combined_005'


def persist_app_state():
    """Try to save any open LibreOffice Impress instance."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.5)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def get_title_underline(slide):
    """
    Return the underline value of the first non-empty run in the title placeholder.
    Returns True if underline is True, or None/False otherwise.
    """
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name.startswith("Title"):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        return run.font.underline
    return None


def get_all_text_content(slide):
    """Return a list of (shape_name, shape_type, top_emu, text) for all text-containing shapes."""
    results = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            full_text = " ".join(
                run.text for para in shape.text_frame.paragraphs
                for run in para.runs if run.text.strip()
            ).strip()
            if full_text:
                results.append((shape.name, shape.shape_type, shape.top, full_text))
    return results


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0

    Scoring:
      Component 1 (0.35): Slide 2 title 'Market Analysis Overview' has underline=True
      Component 2 (0.35): Slide 3 title 'Growth Strategy Roadmap' has underline=True
      Component 3 (0.30): Slide 1 has a text element containing 'CONFIDENTIAL' in the bottom area
                          (top > 70% of slide height)
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slide_height = prs.slide_height  # in EMU

    # Component 1: Slide 2 title has underline formatting (0.35 points)
    try:
        if len(prs.slides) < 2:
            print("FAIL: Component 1 — fewer than 2 slides in presentation")
        else:
            slide2 = prs.slides[1]  # 0-indexed: slide 2
            slide2_title_underline = get_title_underline(slide2)
            if slide2_title_underline is True:
                print(f"PASS: Component 1 — Slide 2 title has underline=True (0.35 pts)")
                total_score += 0.35
            else:
                print(f"FAIL: Component 1 — Slide 2 title underline={slide2_title_underline}, expected True")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 3 title has underline formatting (0.35 points)
    try:
        if len(prs.slides) < 3:
            print("FAIL: Component 2 — fewer than 3 slides in presentation")
        else:
            slide3 = prs.slides[2]  # 0-indexed: slide 3
            slide3_title_underline = get_title_underline(slide3)
            if slide3_title_underline is True:
                print(f"PASS: Component 2 — Slide 3 title has underline=True (0.35 pts)")
                total_score += 0.35
            else:
                print(f"FAIL: Component 2 — Slide 3 title underline={slide3_title_underline}, expected True")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 1 has a 'CONFIDENTIAL' text element in the bottom area (0.30 points)
    # "Bottom area" is defined as top position >= 60% of slide height (permissive to handle slight offsets)
    try:
        slide1 = prs.slides[0]
        bottom_threshold = slide_height * 0.60  # 60% of slide height
        confidential_found = False
        confidential_in_bottom = False

        for shape in slide1.shapes:
            if shape.has_text_frame:
                full_text = " ".join(
                    run.text for para in shape.text_frame.paragraphs
                    for run in para.runs
                ).strip()
                if "CONFIDENTIAL" in full_text.upper():
                    confidential_found = True
                    top_position = shape.top
                    top_ratio = top_position / slide_height
                    if top_position >= bottom_threshold:
                        confidential_in_bottom = True
                        print(f"PASS: Component 3 — 'CONFIDENTIAL' text found at bottom of slide 1 "
                              f"(top={top_position} EMU = {top_ratio:.1%} of slide height, 0.30 pts)")
                    else:
                        print(f"FAIL: Component 3 — 'CONFIDENTIAL' text found at top={top_position} EMU "
                              f"({top_ratio:.1%} of slide height), expected >= 60% of slide height")
                    break

        if not confidential_found:
            print("FAIL: Component 3 — No 'CONFIDENTIAL' text element found on slide 1")

        if confidential_in_bottom:
            total_score += 0.30

    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
