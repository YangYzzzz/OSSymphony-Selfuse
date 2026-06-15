"""
Reward Script: Set text direction to vertical and resize text box on slide 3
Task ID: impress_tct_095
Domain: libreoffice_impress
Scoring:
  Component 1 (0.40) - Text direction set to vertical (stacked top-to-bottom)
  Component 2 (0.35) - Text box resized to tall narrow shape (height > width, ~1.5in x 5in)
  Component 3 (0.25) - Text content 'CREATIVITY' preserved with vertical orientation effective
"""

import os

from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_095'

# Tolerance for dimension checks (relative)
TOLERANCE = 0.15


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]

    # Find the text box on slide 3 (TextBox type, not placeholder)
    textbox = None
    for shape in slide3.shapes:
        if shape.has_text_frame and shape.shape_type == 17:  # TEXT_BOX
            textbox = shape
            break

    if textbox is None:
        # Fallback: find any shape with "CREATIVITY" text
        for shape in slide3.shapes:
            if shape.has_text_frame:
                full_text = "".join(p.text for p in shape.text_frame.paragraphs)
                if "CREATIVITY" in full_text.upper():
                    textbox = shape
                    break

    if textbox is None:
        print("FAIL: No text box found on slide 3")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Found text box '{textbox.name}' on slide 3")
    print(f"INFO: Dimensions: width={textbox.width} EMU ({textbox.width/914400:.2f}in), height={textbox.height} EMU ({textbox.height/914400:.2f}in)")

    # Component 1: Text direction is vertical (0.40 points)
    # The task asks for vertical text (top to bottom, letters stacked)
    # Valid vertical values: 'wordArtVert', 'eaVert', 'vert', 'vert270', 'mongolianVert', 'wordArtVertRtl'
    try:
        bodyPr = textbox.text_frame._txBody.find(qn('a:bodyPr'))
        vert_value = bodyPr.get('vert') if bodyPr is not None else None
        vertical_values = {'wordArtVert', 'eaVert', 'vert', 'vert270', 'mongolianVert', 'wordArtVertRtl'}
        if vert_value is not None and vert_value in vertical_values:
            print(f"PASS: Component 1 -- Text direction is vertical (vert='{vert_value}') (0.40 pts)")
            total_score += 0.40
        else:
            print(f"FAIL: Component 1 -- Text direction not vertical. vert='{vert_value}', expected one of {vertical_values}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Text box is tall and narrow (0.35 points)
    # Golden: ~1.50in wide x ~5.00in tall. Initial was 7.00in wide x 1.20in tall.
    # Key check: height must be significantly greater than width (tall narrow rectangle)
    try:
        w_inches = textbox.width / 914400
        h_inches = textbox.height / 914400

        # Sub-check 2a: Height > Width (the box is tall, not wide) -- 0.20 pts
        if h_inches > w_inches:
            print(f"PASS: Component 2a -- Box is tall (height {h_inches:.2f}in > width {w_inches:.2f}in) (0.20 pts)")
            total_score += 0.20

            # Sub-check 2b: Dimensions approximately match target (~1.5in wide, ~5in tall) -- 0.15 pts
            # Use generous tolerance since the task says "approximately"
            width_ok = w_inches <= 2.5  # narrow enough (was 7.00in originally)
            height_ok = h_inches >= 3.5  # tall enough (was 1.20in originally)
            if width_ok and height_ok:
                print(f"PASS: Component 2b -- Dimensions in target range (w={w_inches:.2f}in <= 2.5, h={h_inches:.2f}in >= 3.5) (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 2b -- Dimensions not in target range (w={w_inches:.2f}in, h={h_inches:.2f}in). Need w<=2.5 and h>=3.5")
        else:
            print(f"FAIL: Component 2 -- Box is NOT tall (height {h_inches:.2f}in <= width {w_inches:.2f}in)")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Text content 'CREATIVITY' is preserved with vertical text effective (0.25 points)
    # We check that the text is still present AND the vertical setting works together with shape
    # This component specifically verifies the combination: text present + vertical + tall shape
    try:
        full_text = ""
        for para in textbox.text_frame.paragraphs:
            full_text += para.text

        has_creativity = "CREATIVITY" in full_text.upper()

        # Get vert again for combined check
        bodyPr = textbox.text_frame._txBody.find(qn('a:bodyPr'))
        vert_value = bodyPr.get('vert') if bodyPr is not None else None
        is_vertical = vert_value is not None and vert_value in {'wordArtVert', 'eaVert', 'vert', 'vert270', 'mongolianVert', 'wordArtVertRtl'}
        is_tall = (textbox.height / 914400) > (textbox.width / 914400)

        if has_creativity and is_vertical and is_tall:
            print(f"PASS: Component 3 -- Text 'CREATIVITY' preserved with vertical+tall layout (0.25 pts)")
            total_score += 0.25
        else:
            reasons = []
            if not has_creativity:
                reasons.append(f"text='{full_text}' missing CREATIVITY")
            if not is_vertical:
                reasons.append(f"not vertical (vert={vert_value})")
            if not is_tall:
                reasons.append("not tall shape")
            print(f"FAIL: Component 3 -- Combined check failed: {', '.join(reasons)}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
