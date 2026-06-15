"""
Initial Setup: Create a presentation with 4 slides where slide 3 has a horizontal text box with 'CREATIVITY'
Task ID: impress_tct_095
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_095'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Creative Design Workshop"
    slide1.placeholders[1].text = "Exploring Visual Communication Techniques"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Workshop Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Morning Session: Typography Fundamentals"
    p2a = body2.add_paragraph()
    p2a.text = "Afternoon Session: Color Theory and Layout"
    p2b = body2.add_paragraph()
    p2b.text = "Evening Session: Hands-on Portfolio Project"

    # --- Slide 3: The key slide with horizontal text box ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add a wide, short text box (horizontal) with "CREATIVITY"
    # Wide: ~7 inches, Short: ~1.2 inches
    left = Inches(1.5)
    top = Inches(3.0)
    width = Inches(7.0)
    height = Inches(1.2)
    txBox = slide3.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "CREATIVITY"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x7A)

    # --- Slide 4: Summary ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Key Takeaways"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Design is not just what it looks like, it is how it works"
    p4a = body4.add_paragraph()
    p4a.text = "Typography shapes the reader's experience and perception"
    p4b = body4.add_paragraph()
    p4b.text = "Experiment with orientation, scale, and direction"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
