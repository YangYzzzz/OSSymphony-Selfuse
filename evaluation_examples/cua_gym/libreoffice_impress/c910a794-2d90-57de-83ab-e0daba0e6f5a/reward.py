"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 109 I accidentally built a bullet list that goes three levels deep. In LibreOffice Impress, how do I promote every level-3 bullet up one step so the slide ends up with only two bullet levels?
Generated: 2025-09-10 16:01:56
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
from pptx import Presentation
from collections import Counter


def verify_promote_bullets(file_path: str, target_slide_index: int = 108) -> float:
    """Reward script for verifying that all level-3 bullets on slide 109 of the
    presentation were promoted one level so that the deepest bullet level is
    now level-2 (pptx level index 1).

    Scoring (progressive):
    • 0.6  – no bullet levels deeper than 2 found (i.e. max level ≤ 1)
    • 0.4  – at least one second-level bullet (level==1) remains, proving
              promotion rather than deletion
    Total possible = 1.0
    """

    max_score = 1.0
    score = 0.0

    # ------------------------------------------------------------------
    # 1. Load the presentation file (no points ‑ prerequisite only)
    # ------------------------------------------------------------------
    print(f"Verifying bullet promotion on slide {target_slide_index + 1} in file: {file_path}\n")

    if not os.path.exists(file_path):
        print("✗ File does not exist")
        return 0.0

    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded successfully with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        return 0.0

    # ------------------------------------------------------------------
    # 2. Validate target slide presence
    # ------------------------------------------------------------------
    if target_slide_index >= len(prs.slides):
        print(f"✗ Slide {target_slide_index + 1} not found in presentation")
        return 0.0

    slide = prs.slides[target_slide_index]

    # ------------------------------------------------------------------
    # 3. Collect paragraph levels on the target slide
    # ------------------------------------------------------------------
    levels = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for p in shape.text_frame.paragraphs:
            if p.text and p.text.strip():
                levels.append(p.level)

    if not levels:
        print("✗ No text paragraphs found on the target slide")
        return 0.0

    level_counter = Counter(levels)
    print("Paragraph level distribution:", level_counter)

    # ------------------------------------------------------------------
    # 4. Scoring – Requirement checks
    # ------------------------------------------------------------------
    # Requirement 1: No bullet levels deeper than 2 (pptx index 1)
    if max(levels) <= 1:
        print("✓ No level-3 bullets (level index 2) present – list depth is 2 or less (0.6 points)")
        score += 0.6
    else:
        print("✗ Found bullet level deeper than 2 – requirement not met (0 points)")

    # Requirement 2: At least one second-level bullet remains to show promotion
    if level_counter.get(1, 0) > 0:
        print("✓ Found at least one second-level bullet (level 1) – indicates promotion rather than removal (0.4 points)")
        score += 0.4
    else:
        print("✗ No second-level bullets found – structure may be incorrect (0 points)")

    # ------------------------------------------------------------------
    final_score = min(score, max_score)
    print(f"\nComputed score: {final_score} / {max_score}")
    print(f"REWARD: {final_score}")
    return final_score


if __name__ == "__main__":
    FILE = "/home/user/on_slide_109_i_accidentally_built_a_bullet_list_that_goes_three_levels_deep_in_libreoffice_impress_h_golden.pptx"
    verify_promote_bullets(FILE)

