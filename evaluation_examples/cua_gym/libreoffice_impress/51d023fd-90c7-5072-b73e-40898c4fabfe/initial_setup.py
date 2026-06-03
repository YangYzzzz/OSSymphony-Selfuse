"""
Initial Setup: Tech Conference presentation with 18 slides, solid white master background.
Task ID: impress_ma_013
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_013'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 alignment=PP_ALIGN.LEFT, color=None):
    """Helper to add a text box with formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()

    # Standard widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # The master slide has a solid white background by default (no changes needed).

    # --- Slide 1: Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "TechForward 2025"
    slide.placeholders[1].text = "Global Innovation Summit\nSan Francisco, CA | September 15-17, 2025"

    # --- Slide 2: Agenda ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conference Agenda"
    body = slide.placeholders[1].text_frame
    body.text = "Day 1: AI & Machine Learning Track"
    body.add_paragraph().text = "Day 2: Cloud Infrastructure & DevOps"
    body.add_paragraph().text = "Day 3: Security & Privacy Workshop"
    body.add_paragraph().text = "Networking Events & Partner Showcases"

    # --- Slide 3: Keynote Speaker ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Keynote: Dr. Elena Vasquez"
    body = slide.placeholders[1].text_frame
    body.text = "Chief AI Officer, Meridian Technologies"
    body.add_paragraph().text = ""
    body.add_paragraph().text = '"The Next Frontier: Autonomous Systems in Enterprise"'
    body.add_paragraph().text = "15+ years in AI research"
    body.add_paragraph().text = "Former Lead Scientist at DeepCore Labs"

    # --- Slide 4: AI Track Overview ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "AI & Machine Learning Track"
    body = slide.placeholders[1].text_frame
    body.text = "Foundation Models: From Research to Production"
    body.add_paragraph().text = "Responsible AI: Bias Detection & Mitigation"
    body.add_paragraph().text = "Real-time Inference at Scale"
    body.add_paragraph().text = "Multi-modal AI Applications"
    body.add_paragraph().text = "Panel: AI Regulation & Compliance"

    # --- Slide 5: Cloud Track Overview ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Cloud Infrastructure & DevOps"
    body = slide.placeholders[1].text_frame
    body.text = "Kubernetes at Scale: Lessons from Production"
    body.add_paragraph().text = "Serverless 2.0: Beyond Functions"
    body.add_paragraph().text = "Infrastructure as Code Best Practices"
    body.add_paragraph().text = "Cost Optimization Strategies"

    # --- Slide 6: Security Track ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Security & Privacy Workshop"
    body = slide.placeholders[1].text_frame
    body.text = "Zero Trust Architecture Implementation"
    body.add_paragraph().text = "Supply Chain Security"
    body.add_paragraph().text = "Post-Quantum Cryptography Readiness"
    body.add_paragraph().text = "GDPR & Global Privacy Updates"

    # --- Slide 7: Speaker - Marcus Chen ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Marcus Chen"
    body = slide.placeholders[1].text_frame
    body.text = "VP of Engineering, CloudScale Inc."
    body.add_paragraph().text = ""
    body.add_paragraph().text = "Topic: Building Resilient Distributed Systems"
    body.add_paragraph().text = "Author of 'Microservices at Scale'"
    body.add_paragraph().text = "Previously: Principal Engineer at AmazonTech"

    # --- Slide 8: Speaker - Sarah Okafor ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Dr. Sarah Okafor"
    body = slide.placeholders[1].text_frame
    body.text = "Director of AI Ethics, NovaTech"
    body.add_paragraph().text = ""
    body.add_paragraph().text = "Topic: Fairness in Foundation Models"
    body.add_paragraph().text = "IEEE Ethics Committee Member"
    body.add_paragraph().text = "Ph.D. Stanford University"

    # --- Slide 9: Industry Stats ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Industry Growth Metrics"
    body = slide.placeholders[1].text_frame
    body.text = "Global AI market: $407B by 2027 (Gartner)"
    body.add_paragraph().text = "Cloud spending: $680B in 2024"
    body.add_paragraph().text = "Cybersecurity market: $298B projected"
    body.add_paragraph().text = "DevOps adoption: 83% of enterprises"

    # --- Slide 10: Workshop Details ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Hands-On Workshops"
    body = slide.placeholders[1].text_frame
    body.text = "Workshop A: Building RAG Pipelines (3 hours)"
    body.add_paragraph().text = "Workshop B: Kubernetes Security Hardening (2 hours)"
    body.add_paragraph().text = "Workshop C: LLM Fine-tuning Lab (4 hours)"
    body.add_paragraph().text = "Workshop D: Observability & Monitoring (2 hours)"

    # --- Slide 11: Networking Events ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Networking & Social Events"
    body = slide.placeholders[1].text_frame
    body.text = "Welcome Reception - Rooftop Terrace (Day 1, 6PM)"
    body.add_paragraph().text = "Startup Pitch Competition (Day 2, 2PM)"
    body.add_paragraph().text = "Partner Showcase Exhibition Hall (All Days)"
    body.add_paragraph().text = "Closing Gala - The Embarcadero (Day 3, 7PM)"

    # --- Slide 12: Sponsor Spotlight ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Platinum Sponsors"
    body = slide.placeholders[1].text_frame
    body.text = "Meridian Technologies - AI Infrastructure"
    body.add_paragraph().text = "CloudScale Inc. - Enterprise Cloud"
    body.add_paragraph().text = "NovaTech - AI Ethics & Governance"
    body.add_paragraph().text = "CyberShield Corp. - Security Solutions"

    # --- Slide 13: Venue Information ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Venue: Moscone Center"
    body = slide.placeholders[1].text_frame
    body.text = "747 Howard Street, San Francisco, CA 94103"
    body.add_paragraph().text = "North Hall: Main Stage & Keynotes"
    body.add_paragraph().text = "South Hall: Workshops & Labs"
    body.add_paragraph().text = "West Hall: Exhibition & Partner Booths"

    # --- Slide 14: Schedule Day 1 ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Day 1 Schedule - September 15"
    body = slide.placeholders[1].text_frame
    body.text = "8:00 AM - Registration & Breakfast"
    body.add_paragraph().text = "9:00 AM - Opening Keynote: Dr. Vasquez"
    body.add_paragraph().text = "10:30 AM - AI Track Sessions"
    body.add_paragraph().text = "12:00 PM - Lunch & Networking"
    body.add_paragraph().text = "1:30 PM - Cloud Track Sessions"
    body.add_paragraph().text = "4:00 PM - Lightning Talks"
    body.add_paragraph().text = "6:00 PM - Welcome Reception"

    # --- Slide 15: Schedule Day 2 ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Day 2 Schedule - September 16"
    body = slide.placeholders[1].text_frame
    body.text = "8:30 AM - Breakfast & Networking"
    body.add_paragraph().text = "9:00 AM - Security Keynote"
    body.add_paragraph().text = "10:30 AM - Hands-On Workshops"
    body.add_paragraph().text = "12:30 PM - Lunch"
    body.add_paragraph().text = "2:00 PM - Startup Pitch Competition"
    body.add_paragraph().text = "4:30 PM - Panel Discussions"

    # --- Slide 16: Schedule Day 3 ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Day 3 Schedule - September 17"
    body = slide.placeholders[1].text_frame
    body.text = "8:30 AM - Final Day Breakfast"
    body.add_paragraph().text = "9:00 AM - Deep Dive Sessions"
    body.add_paragraph().text = "11:00 AM - Closing Keynote"
    body.add_paragraph().text = "12:30 PM - Lunch"
    body.add_paragraph().text = "2:00 PM - Award Ceremony"
    body.add_paragraph().text = "7:00 PM - Closing Gala"

    # --- Slide 17: Call to Action ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Register Now"
    body = slide.placeholders[1].text_frame
    body.text = "Early Bird: $599 (until July 31)"
    body.add_paragraph().text = "Standard: $899"
    body.add_paragraph().text = "VIP All-Access: $1,499"
    body.add_paragraph().text = "Group Discount: 15% off for 5+ attendees"
    body.add_paragraph().text = ""
    body.add_paragraph().text = "www.techforward2025.com/register"

    # --- Slide 18: Thank You / Contact ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Thank You"
    body = slide.placeholders[1].text_frame
    body.text = "info@techforward2025.com"
    body.add_paragraph().text = "@TechForward2025"
    body.add_paragraph().text = "#TF2025"
    body.add_paragraph().text = ""
    body.add_paragraph().text = "See you in San Francisco!"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
