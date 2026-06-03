"""
Initial Setup: Configure slide show settings presentation
Task ID: impress_gf3_050
Domain: libreoffice_impress

Creates a 14-slide presentation with default slide show settings.
The slide show starts from slide 1, uses default cursor, no black end slide.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_050'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q4 Strategic Planning"
    slide1.placeholders[1].text = "Meridian Technologies Inc.\nOctober 2025"

    # Slide 2: Agenda
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    tf = slide2.placeholders[1].text_frame
    tf.text = "1. Executive Summary"
    for item in ["2. Revenue Analysis", "3. Market Expansion Plans",
                  "4. Product Roadmap", "5. Team & Hiring",
                  "6. Budget Allocations", "7. Q&A"]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

    # Slide 3: Executive Summary
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Executive Summary"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Total revenue grew 18% YoY to $42.7M"
    for line in ["Customer acquisition cost reduced by 12%",
                 "Net promoter score improved from 67 to 74",
                 "Three new enterprise clients onboarded in Q3"]:
        p = tf3.add_paragraph()
        p.text = line

    # Slide 4: Revenue Breakdown
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Revenue Breakdown"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "SaaS Subscriptions: $28.3M (+22%)"
    for line in ["Professional Services: $8.9M (+11%)",
                 "Hardware Licensing: $3.2M (-5%)",
                 "Support Contracts: $2.3M (+8%)"]:
        p = tf4.add_paragraph()
        p.text = line

    # Slide 5: Market Analysis
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Market Analysis"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "TAM expanded to $4.2B in cloud infrastructure"
    for line in ["Competitor consolidation creating opportunities",
                 "APAC region showing 35% growth potential",
                 "EU regulatory compliance as differentiator"]:
        p = tf5.add_paragraph()
        p.text = line

    # Slide 6: Product Roadmap
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Product Roadmap - H1 2026"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "AI-powered analytics dashboard (Jan 2026)"
    for line in ["Multi-cloud orchestration layer (Feb 2026)",
                 "Enhanced security compliance module (Mar 2026)",
                 "Mobile app v3.0 release (Apr 2026)"]:
        p = tf6.add_paragraph()
        p.text = line

    # Slide 7: Engineering Team
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Engineering Team Overview"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "Current headcount: 87 engineers"
    for line in ["Frontend: 24 | Backend: 31 | DevOps: 18 | QA: 14",
                 "Open positions: 12 (6 backend, 4 ML, 2 SRE)",
                 "Average tenure: 2.4 years",
                 "Attrition rate: 8.2% (industry avg: 13.5%)"]:
        p = tf7.add_paragraph()
        p.text = line

    # Slide 8: Customer Success
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Customer Success Metrics"
    tf8 = slide8.placeholders[1].text_frame
    tf8.text = "Enterprise clients: 47 (+8 in Q3)"
    for line in ["Monthly active users: 234,500",
                 "Support ticket resolution: 4.2 hrs avg",
                 "Customer satisfaction: 4.6/5.0",
                 "Renewal rate: 94.3%"]:
        p = tf8.add_paragraph()
        p.text = line

    # Slide 9: Financial Overview
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Financial Overview"
    tf9 = slide9.placeholders[1].text_frame
    tf9.text = "Gross margin: 72.4% (up from 68.1%)"
    for line in ["Operating expenses: $31.2M",
                 "EBITDA: $11.5M (27% margin)",
                 "Cash reserves: $18.7M",
                 "Burn rate reduced 15% QoQ"]:
        p = tf9.add_paragraph()
        p.text = line

    # Slide 10: Partnerships
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "Strategic Partnerships"
    tf10 = slide10.placeholders[1].text_frame
    tf10.text = "AWS Advanced Partner status achieved"
    for line in ["Microsoft Azure co-sell agreement signed",
                 "Integration with Salesforce marketplace",
                 "Joint venture with DataBridge Analytics"]:
        p = tf10.add_paragraph()
        p.text = line

    # Slide 11: Risk Assessment
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    slide11.shapes.title.text = "Risk Assessment"
    tf11 = slide11.placeholders[1].text_frame
    tf11.text = "Supply chain disruption: Medium risk"
    for line in ["Talent acquisition in ML/AI: High risk",
                 "Regulatory changes (EU AI Act): Medium risk",
                 "Currency fluctuation impact: Low risk",
                 "Cybersecurity threat landscape: Medium risk"]:
        p = tf11.add_paragraph()
        p.text = line

    # Slide 12: Budget Allocation
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])
    slide12.shapes.title.text = "Q4 Budget Allocation"
    tf12 = slide12.placeholders[1].text_frame
    tf12.text = "R&D: $14.2M (45%)"
    for line in ["Sales & Marketing: $8.5M (27%)",
                 "Operations: $4.7M (15%)",
                 "G&A: $2.8M (9%)",
                 "Contingency: $1.2M (4%)"]:
        p = tf12.add_paragraph()
        p.text = line

    # Slide 13: Next Steps
    slide13 = prs.slides.add_slide(prs.slide_layouts[1])
    slide13.shapes.title.text = "Next Steps & Action Items"
    tf13 = slide13.placeholders[1].text_frame
    tf13.text = "Finalize 2026 hiring plan by Nov 15"
    for line in ["Complete AWS migration Phase 2 by Dec 1",
                 "Launch beta for AI analytics by Jan 15",
                 "Board presentation scheduled for Dec 18",
                 "All-hands meeting on Oct 28"]:
        p = tf13.add_paragraph()
        p.text = line

    # Slide 14: Thank You / Q&A
    slide14 = prs.slides.add_slide(prs.slide_layouts[0])
    slide14.shapes.title.text = "Thank You"
    slide14.placeholders[1].text = "Questions & Discussion\ncontact@meridiantech.com"

    # Save with default slide show settings (no custom showPr element)
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
