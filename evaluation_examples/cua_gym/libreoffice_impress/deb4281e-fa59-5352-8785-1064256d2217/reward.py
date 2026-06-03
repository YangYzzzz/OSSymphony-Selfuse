"""
FINAL REWARD SCRIPT - SUCCESS
Task: Insert image 'Desktop/1.png' anchored As Character with width 7.5 cm (keep ratio).
Generated: 2025-10-17 05:54:16
Status: success
Model: azure-o3
Total Steps: 8
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def verify_insert_image_task(file_path: str) -> float:
    """Verify that the presentation contains an image anchored ‘As Character’
    (i.e., a Picture shape) that has a width of 7.5 cm (±≈0.05 cm) while
    preserving the original aspect ratio.

    Scoring (progressive):
      • 0.4 – A picture is present in the presentation.
      • 0.3 – That picture’s width is within tolerance of 7.5 cm.
      • 0.3 – The picture’s aspect ratio is preserved (≤3 % deviation).
    Returns a float between 0.0 and 1.0 (inclusive).
    """

    # Constants for conversion and tolerances
    CM_TO_EMU = 360_000            # 1 cm = 360 000 EMU
    EXPECTED_WIDTH_CM = 7.5
    EXPECTED_WIDTH_EMU = EXPECTED_WIDTH_CM * CM_TO_EMU
    WIDTH_TOLERANCE_EMU = 20_000   # ≈0.055 cm
    RATIO_TOLERANCE = 0.03         # 3 %

    print(f"Verifying file: {file_path}")

    # 0. Check file existence
    if not os.path.exists(file_path):
        print("✗ File not found.")
        return 0.0

    # 1. Attempt to load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        return 0.0

    # Flags for progressive scoring
    picture_found = False
    width_ok = False
    ratio_ok = False

    # 2. Iterate over all shapes in all slides
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_found = True
                print(f"✓ Picture found (slide {slide_idx}, shape {shape_idx})")
                print(f"  Shape size: width={shape.width} EMU, height={shape.height} EMU")

                # 2a. Width verification
                if abs(shape.width - EXPECTED_WIDTH_EMU) <= WIDTH_TOLERANCE_EMU:
                    width_ok = True
                    print(f"  ✓ Width within tolerance of 7.5 cm (±0.055 cm)")
                else:
                    print(f"  ✗ Width {shape.width} EMU outside tolerance")

                # 2b. Aspect-ratio verification
                try:
                    # Shape aspect ratio (EMU are proportional, so division is safe)
                    shape_ratio = shape.width / shape.height if shape.height else None
                    # Original image pixel ratio (using python-pptx Pillow size tuple)
                    img_w, img_h = shape.image.size  # (px, px)
                    image_ratio = img_w / img_h if img_h else None

                    if shape_ratio and image_ratio:
                        diff = abs(shape_ratio - image_ratio) / image_ratio
                        if diff <= RATIO_TOLERANCE:
                            ratio_ok = True
                            print(f"  ✓ Aspect ratio preserved (diff {diff:.2%})")
                        else:
                            print(f"  ✗ Aspect ratio changed (diff {diff:.2%})")
                except Exception as e:
                    print(f"  ✗ Error checking aspect ratio: {e}")

    # 3. Progressive scoring based on verification results
    total_score = 0.0
    if picture_found:
        total_score += 0.4
    if picture_found and width_ok:
        total_score += 0.3
    if picture_found and width_ok and ratio_ok:
        total_score += 0.3

    print(f"Total score: {total_score}")
    return min(total_score, 1.0)


# -------------------- MAIN EXECUTION --------------------
if __name__ == "__main__":
    # Path provided by the grading runner / task definition
    FILE_PATH = "/home/user/insert_image_desktop1png_anchored_as_character_with_width_75_cm_keep_ratio.pptx"

    reward_value = verify_insert_image_task(FILE_PATH)
    print(f"REWARD: {reward_value}")
