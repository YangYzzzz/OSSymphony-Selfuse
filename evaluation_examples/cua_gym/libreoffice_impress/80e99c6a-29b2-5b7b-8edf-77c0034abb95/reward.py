"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 177, could you switch the title text color to #FF9900 (the “Orange 2” swatch) and enable the shadow effect with a 0.1 cm offset? Thanks!
Generated: 2025-09-10 17:58:39
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from lxml import etree

# ----------------------------------------------------------------------------
# Reward Script for Verifying Task Completion
# Task: On slide 177, switch the title text color to #FF9900 (Orange 2)
#       and enable the shadow effect with a 0.1 cm offset.
# ----------------------------------------------------------------------------
# Scoring
#   • 0.5 points – Title text colour on slide 177 is exactly #FF9900
#   • 0.5 points – Title text has an outer shadow with a distance of 0.1 cm
#   • 1.0 points – Both requirements satisfied
# Progressive scoring is used; partial completion is rewarded.
# ----------------------------------------------------------------------------

FILE_PATH = "/home/user/on_slide_177_could_you_switch_the_title_text_color_to_ff9900_the_orange_2_swatch_and_enable_the_shad_golden.pptx"

# Helper ‑ locate the title shape (prefer placeholder titles, otherwise first text shape)
def _find_title_shape(slide):
    # 1) Placeholder title types first
    for shape in slide.shapes:
        try:
            if shape.is_placeholder and shape.has_text_frame:
                if shape.placeholder_format.type in (0, 1, 5, 6, 7, 9, 11):  # title-related types
                    return shape
        except Exception:
            pass
    # 2) Fallback – first shape with text
    for shape in slide.shapes:
        if shape.has_text_frame:
            return shape
    return None

# Check colour of every run inside the title shape
def _verify_title_colour(shape, expected_rgb):
    if shape is None or not shape.has_text_frame:
        print("✗ No title shape with a text-frame found on slide 177")
        return False

    every_run_correct = True
    any_run_correct = False

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run_rgb = run.font.color.rgb
            if run_rgb == expected_rgb:
                any_run_correct = True
            else:
                every_run_correct = False
            print(f"  • Run colour = {run_rgb}")

    if any_run_correct:
        if every_run_correct:
            print("✓ Title text colour: all runs are #FF9900")
        else:
            print("✓ Title text colour: at least one run is #FF9900 (partial)")
        return True
    else:
        print("✗ Title text colour #FF9900 not found in any run")
        return False

# Check for outer shadow on the title shape
# 0.1 cm = 0.1 * 360 000 EMU = 36 000 EMU
def _verify_shadow(shape, expected_dist_emu):
    if shape is None:
        print("✗ Cannot verify shadow – no title shape found")
        return False

    xml_str = shape._element.xml
    root = etree.fromstring(xml_str.encode())
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    shadows = root.xpath(".//a:outerShdw", namespaces=ns)

    if not shadows:
        print("✗ No outer shadow found on the title text")
        return False

    for sh in shadows:
        dist_raw = sh.get("dist")
        try:
            dist_val = int(dist_raw)
            print(f"  • Detected shadow distance = {dist_val} EMU")
            # Accept ±1000 EMU (~0.003 cm) tolerance for rounding
            if abs(dist_val - expected_dist_emu) <= 1000:
                print("✓ Shadow distance matches the expected 0.1 cm")
                return True
        except (TypeError, ValueError):
            pass

    print("✗ Shadow found, but distance does not match 0.1 cm")
    return False

# Main verification routine
def verify_task(file_path):
    max_score = 1.0
    score = 0.0

    if not os.path.exists(file_path):
        print("✗ Presentation file not found:", file_path)
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print("✗ Failed to load presentation:", e)
        return 0.0

    if len(prs.slides) < 177:
        print(f"✗ Presentation only has {len(prs.slides)} slides – slide 177 is missing")
        return 0.0

    slide177 = prs.slides[176]  # zero-based index
    title_shape = _find_title_shape(slide177)

    # Requirement 1: text colour
    expected_rgb = RGBColor(255, 153, 0)  # #FF9900
    if _verify_title_colour(title_shape, expected_rgb):
        score += 0.5

    # Requirement 2: shadow distance
    expected_dist_emu = int(round(0.1 * 360000))  # 36 000 EMU
    if _verify_shadow(title_shape, expected_dist_emu):
        score += 0.5

    final_score = min(score, max_score)
    print(f"Total Score = {final_score}/{max_score}")
    return final_score

if __name__ == "__main__":
    reward = verify_task(FILE_PATH)
    print(f"REWARD: {reward}")
