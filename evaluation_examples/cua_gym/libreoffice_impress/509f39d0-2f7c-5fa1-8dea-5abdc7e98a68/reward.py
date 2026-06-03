"""
Reward Script: Student feedback form with star rating on slide 9
Task ID: impress_teach_090
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): 'Rate this lecture:' text box on slide 9
  Component 2 (0.25): Exactly 5 five-point star shapes on slide 9
  Component 3 (0.30): Stars 1-3 filled gold (#FFD700)
  Component 4 (0.20): Stars 4-5 filled light gray (#E0E0E0)
"""

import os

from pptx import Presentation
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_090'


def persist_app_state(domain):
    """Save any unsaved GUI edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            import time
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_star_shapes(slide):
    """Return list of shapes that are 5-point stars (prst='star5'), sorted by left position."""
    stars = []
    for shape in slide.shapes:
        if shape.shape_type == 1:  # AUTO_SHAPE
            sp = shape._element
            prstGeom = sp.find('.//' + qn('a:prstGeom'))
            if prstGeom is not None and prstGeom.get('prst') == 'star5':
                stars.append(shape)
    # Sort by horizontal position (left to right)
    stars.sort(key=lambda s: s.left)
    return stars


def get_fill_color_hex(shape):
    """Get the solid fill color of a shape as uppercase hex string, or None."""
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb).upper()
    except Exception:
        pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 9 slides
    if len(prs.slides) < 9:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 9")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[8]  # 0-indexed, slide 9

    # Component 1: 'Rate this lecture:' text box on slide 9 (0.25 points)
    try:
        found_rate_text = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                full_text = shape.text_frame.text.strip()
                if 'rate this lecture' in full_text.lower():
                    found_rate_text = True
                    break
        if found_rate_text:
            print(f"PASS: Component 1 - 'Rate this lecture:' text found on slide 9 (0.25 pts)")
            total_score += 0.25
        else:
            texts = [s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame]
            print(f"FAIL: Component 1 - 'Rate this lecture:' not found. Texts on slide 9: {texts}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Exactly 5 five-point star shapes on slide 9 (0.25 points)
    try:
        stars = get_star_shapes(slide)
        num_stars = len(stars)
        if num_stars == 5:
            print(f"PASS: Component 2 - Found exactly 5 five-point star shapes (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 - Expected 5 five-point stars, found {num_stars}")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Stars 1-3 filled gold #FFD700 (0.30 points)
    try:
        stars = get_star_shapes(slide)
        if len(stars) >= 3:
            gold_count = 0
            for i in range(3):
                color = get_fill_color_hex(stars[i])
                if color == 'FFD700':
                    gold_count += 1
                else:
                    print(f"  Star {i+1} fill: {color} (expected FFD700)")
            if gold_count == 3:
                print(f"PASS: Component 3 - Stars 1-3 all have gold fill #FFD700 (0.30 pts)")
                total_score += 0.30
            else:
                print(f"FAIL: Component 3 - Only {gold_count}/3 stars have gold fill")
        else:
            print(f"FAIL: Component 3 - Not enough star shapes to check (need 3, have {len(stars)})")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Stars 4-5 filled light gray #E0E0E0 (0.20 points)
    try:
        stars = get_star_shapes(slide)
        if len(stars) >= 5:
            gray_count = 0
            for i in range(3, 5):
                color = get_fill_color_hex(stars[i])
                if color == 'E0E0E0':
                    gray_count += 1
                else:
                    print(f"  Star {i+1} fill: {color} (expected E0E0E0)")
            if gray_count == 2:
                print(f"PASS: Component 4 - Stars 4-5 both have gray fill #E0E0E0 (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 4 - Only {gray_count}/2 stars have gray fill")
        else:
            print(f"FAIL: Component 4 - Not enough star shapes to check (need 5, have {len(stars)})")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
