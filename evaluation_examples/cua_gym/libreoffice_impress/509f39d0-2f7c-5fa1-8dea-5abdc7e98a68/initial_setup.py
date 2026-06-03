"""
Initial Setup: Create a 10-slide teaching presentation with slide 9 titled 'Your Feedback Matters'
Task ID: impress_teach_090
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_090'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_text):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.text = body_text
    return slide


def add_blank_slide_with_title(prs, title_text):
    """Add a slide using Title Only layout (index 5) with just a title."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8.5), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x7A)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(prs, "Introduction to Data Science", "Fall 2025 - Dr. Sarah Chen")

    # Slide 2: Course Overview
    add_content_slide(prs, "Course Overview", (
        "This course covers fundamental concepts of data science including:\n"
        "- Statistical analysis and probability\n"
        "- Machine learning algorithms\n"
        "- Data visualization techniques\n"
        "- Real-world case studies"
    ))

    # Slide 3: Learning Objectives
    add_content_slide(prs, "Learning Objectives", (
        "By the end of this course, students will be able to:\n"
        "1. Apply statistical methods to real datasets\n"
        "2. Build and evaluate predictive models\n"
        "3. Create effective data visualizations\n"
        "4. Present analytical findings clearly"
    ))

    # Slide 4: Key Topics
    add_content_slide(prs, "Key Topics for This Week", (
        "Week 7 - Regression Analysis\n"
        "- Simple linear regression\n"
        "- Multiple regression models\n"
        "- Evaluating model fit (R-squared, RMSE)\n"
        "- Practical exercise with housing price data"
    ))

    # Slide 5: Data Visualization
    add_content_slide(prs, "Data Visualization Best Practices", (
        "Effective charts and graphs should:\n"
        "- Use appropriate chart types for the data\n"
        "- Minimize chartjunk and unnecessary decoration\n"
        "- Label axes clearly with units\n"
        "- Choose accessible color palettes"
    ))

    # Slide 6: Case Study
    add_content_slide(prs, "Case Study: Urban Transportation Analysis", (
        "Dataset: 15,000 ride records from Metro Transit Authority\n"
        "Objective: Predict peak demand patterns\n"
        "Methods: Time series decomposition, ARIMA modeling\n"
        "Result: 23% improvement in resource allocation"
    ))

    # Slide 7: Group Project Guidelines
    add_content_slide(prs, "Group Project Guidelines", (
        "Teams of 3-4 students\n"
        "- Choose a dataset of at least 5,000 records\n"
        "- Apply at least two analytical techniques\n"
        "- Prepare a 15-minute presentation\n"
        "- Submit written report by December 12, 2025"
    ))

    # Slide 8: Resources
    add_content_slide(prs, "Recommended Resources", (
        "Textbook: 'Python for Data Analysis' by Wes McKinney\n"
        "Online: Kaggle datasets and tutorials\n"
        "Software: Jupyter Notebook, scikit-learn, pandas\n"
        "Office hours: Tuesdays 2-4 PM, Room 312B"
    ))

    # Slide 9: Feedback slide - title only, NO star shapes, NO rating text
    add_blank_slide_with_title(prs, "Your Feedback Matters")

    # Slide 10: Thank You
    add_title_slide(prs, "Thank You!", "Questions? Email: s.chen@university.edu")

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Slide count: {len(prs.slides)}')

    # Launch in Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
