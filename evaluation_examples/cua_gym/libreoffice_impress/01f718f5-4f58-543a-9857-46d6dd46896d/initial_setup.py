"""
Initial Setup: Create a presentation with a line chart on slide 3 (no chart/axis titles)
Task ID: impress_tct_042
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_042'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Growth Analysis Report"
    slide1.placeholders[1].text = "Fiscal Year Performance Overview\nPrepared by: Analytics Division"

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "This report presents the annual growth rate analysis across multiple fiscal years."
    p2 = body2.add_paragraph()
    p2.text = "Key highlights include sustained growth from 2018 through 2023, with a notable acceleration in recent years driven by digital transformation initiatives."
    p3 = body2.add_paragraph()
    p3.text = "The following slides contain detailed charts and breakdowns of growth metrics by business unit."

    # --- Slide 3: Line Chart (NO title, NO axis titles) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Add a title text box manually
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Annual Growth Rate"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Chart data
    chart_data = CategoryChartData()
    chart_data.categories = ['2018', '2019', '2020', '2021', '2022', '2023']
    chart_data.add_series('Revenue Growth', (3.2, 4.1, 2.8, 5.6, 7.3, 8.1))
    chart_data.add_series('Profit Growth', (2.1, 3.5, 1.9, 4.2, 6.0, 7.4))

    # Add chart - line chart
    chart_shape = slide3.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(1.0), Inches(1.2), Inches(10), Inches(5.5),
        chart_data,
    )
    chart = chart_shape.chart

    # Explicitly disable chart title
    chart.has_title = False

    # Ensure no axis titles
    category_axis = chart.category_axis
    category_axis.has_title = False

    value_axis = chart.value_axis
    value_axis.has_title = False

    # Style the chart lines
    plot = chart.plots[0]
    plot.smooth = False

    series0 = chart.series[0]
    series0.format.line.color.rgb = RGBColor(0x44, 0x72, 0xC4)
    series0.format.line.width = Pt(2.5)

    series1 = chart.series[1]
    series1.format.line.color.rgb = RGBColor(0xED, 0x7D, 0x31)
    series1.format.line.width = Pt(2.5)

    # --- Slide 4: Regional Breakdown ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Regional Performance"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "North America: Led growth at 9.2% in 2023, driven by strong enterprise adoption."
    p4a = body4.add_paragraph()
    p4a.text = "Europe: Steady recovery post-2020, reaching 6.8% growth in 2023."
    p4b = body4.add_paragraph()
    p4b.text = "Asia-Pacific: Fastest growing region at 11.3%, fueled by emerging market expansion."
    p4c = body4.add_paragraph()
    p4c.text = "Latin America: Moderate growth of 4.5%, with infrastructure investments expected to accelerate."

    # --- Slide 5: Conclusion ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Outlook & Recommendations"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Based on current trends, projected growth for 2024 is estimated at 8.5-9.0%."
    p5a = body5.add_paragraph()
    p5a.text = "Key recommendations: Increase investment in Asia-Pacific markets and accelerate digital product launches."
    p5b = body5.add_paragraph()
    p5b.text = "Risk factors include regulatory changes in the EU and supply chain disruptions in emerging markets."

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
