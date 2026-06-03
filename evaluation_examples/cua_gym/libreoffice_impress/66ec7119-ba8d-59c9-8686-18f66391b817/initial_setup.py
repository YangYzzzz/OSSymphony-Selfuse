"""
Initial Setup: Create a 6-slide history lecture presentation
Task ID: impstruct_007
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impstruct_007'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    """Add a title-layout slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points):
    """Add a title+content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            body.paragraphs[0].text = point
        else:
            p = body.add_paragraph()
            p.text = point
    return slide


def create_initial():
    prs = Presentation()

    # Slide 1: Title
    add_title_slide(
        prs,
        "The Rise and Fall of Ancient Civilizations",
        "A Comprehensive Survey — Prof. Elena Vasquez\nWorld History 301 — Spring 2025"
    )

    # Slide 2: Era 1
    add_content_slide(prs, "Era 1: Mesopotamia (3500–2000 BCE)", [
        "Sumerian city-states established along the Tigris and Euphrates rivers",
        "Invention of cuneiform writing around 3200 BCE enabled record-keeping",
        "The Code of Ur-Nammu (c. 2100 BCE) — earliest known legal code",
        "Ziggurats served as religious and administrative centers",
        "Advanced irrigation systems supported agriculture in arid regions",
    ])

    # Slide 3: Era 2
    add_content_slide(prs, "Era 2: Classical Greece (800–146 BCE)", [
        "Athens developed direct democracy under Cleisthenes (508 BCE)",
        "The Persian Wars (499–449 BCE) united Greek city-states temporarily",
        "Golden Age of Pericles saw construction of the Parthenon (447–432 BCE)",
        "Philosophical traditions founded by Socrates, Plato, and Aristotle",
        "Alexander the Great spread Hellenistic culture across three continents",
    ])

    # Slide 4: Era 3
    add_content_slide(prs, "Era 3: The Roman Empire (27 BCE–476 CE)", [
        "Augustus Caesar established the Principate in 27 BCE",
        "Pax Romana (27 BCE–180 CE) — two centuries of relative stability",
        "Engineering marvels: aqueducts, roads, and the Colosseum",
        "Adoption of Christianity as state religion under Theodosius I (380 CE)",
        "Fall of the Western Empire in 476 CE due to invasions and internal decay",
    ])

    # Slide 5: Era 4
    add_content_slide(prs, "Era 4: The Medieval Period (500–1500 CE)", [
        "Feudal system organized society into lords, vassals, and serfs",
        "The Byzantine Empire preserved Roman law and Greek scholarship",
        "Charlemagne crowned Holy Roman Emperor in 800 CE",
        "The Crusades (1096–1291) reshaped trade routes and cultural exchange",
        "The Black Death (1347–1351) killed roughly one-third of Europe's population",
    ])

    # Slide 6: Summary
    add_content_slide(prs, "Summary: Patterns Across Civilizations", [
        "Agricultural surplus enabled urbanization and specialization",
        "Legal codification strengthened governance and social order",
        "Military expansion spread culture but often led to overextension",
        "Disease and climate shifts repeatedly disrupted established powers",
        "Each era's innovations built foundations for subsequent civilizations",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
