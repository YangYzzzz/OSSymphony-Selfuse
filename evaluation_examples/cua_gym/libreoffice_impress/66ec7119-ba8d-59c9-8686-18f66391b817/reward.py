"""
Reward Script: Reverse the order of slides 2 through 5
Task ID: impstruct_007
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.25): Slide 2 title contains "Era 4"
  - Component 2 (0.25): Slide 3 title contains "Era 3"
  - Component 3 (0.25): Slide 4 title contains "Era 2"
  - Component 4 (0.25): Slide 5 title contains "Era 1"
  Gate: Total slides == 6, Slide 1 unchanged, Slide 6 unchanged
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impstruct_007'


def get_slide_title(slide):
    """Extract first non-empty text from a slide's shapes."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    return text
    return ""


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Gate: Must have exactly 6 slides
    num_slides = len(prs.slides)
    if num_slides != 6:
        print(f"GATE FAIL: Expected 6 slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)

    # Gate: Slide 1 must still be the title slide (unchanged)
    slide1_title = get_slide_title(slides[0])
    if "Rise and Fall" not in slide1_title and "Ancient Civilizations" not in slide1_title:
        print(f"GATE FAIL: Slide 1 title appears changed: {slide1_title!r}")
        print("REWARD: 0.0")
        return 0.0

    # Gate: Slide 6 must still be the summary slide (unchanged)
    slide6_title = get_slide_title(slides[5])
    if "Summary" not in slide6_title:
        print(f"GATE FAIL: Slide 6 title appears changed: {slide6_title!r}")
        print("REWARD: 0.0")
        return 0.0

    print(f"GATE PASS: 6 slides, Slide 1 = title, Slide 6 = summary")

    # Component 1: Slide 2 should now contain "Era 4" (0.25 points)
    # In initial, slide 2 has "Era 1" — this checks the reversal
    try:
        slide2_title = get_slide_title(slides[1])
        if "Era 4" in slide2_title:
            print(f"PASS: Component 1 — Slide 2 title contains 'Era 4': {slide2_title!r} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — Expected 'Era 4' in slide 2, found: {slide2_title!r}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 3 should now contain "Era 3" (0.25 points)
    # In initial, slide 3 has "Era 2" — this checks the reversal
    try:
        slide3_title = get_slide_title(slides[2])
        if "Era 3" in slide3_title:
            print(f"PASS: Component 2 — Slide 3 title contains 'Era 3': {slide3_title!r} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 — Expected 'Era 3' in slide 3, found: {slide3_title!r}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 4 should now contain "Era 2" (0.25 points)
    # In initial, slide 4 has "Era 3" — this checks the reversal
    try:
        slide4_title = get_slide_title(slides[3])
        if "Era 2" in slide4_title:
            print(f"PASS: Component 3 — Slide 4 title contains 'Era 2': {slide4_title!r} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 — Expected 'Era 2' in slide 4, found: {slide4_title!r}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 5 should now contain "Era 1" (0.25 points)
    # In initial, slide 5 has "Era 4" — this checks the reversal
    try:
        slide5_title = get_slide_title(slides[4])
        if "Era 1" in slide5_title:
            print(f"PASS: Component 4 — Slide 5 title contains 'Era 1': {slide5_title!r} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 4 — Expected 'Era 1' in slide 5, found: {slide5_title!r}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook: save any unsaved LibreOffice state
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
