"""
FINAL REWARD SCRIPT - SUCCESS
Task: Almost done editing, but the marketing team needs the deck in native PowerPoint format. In LibreOffice Impress, what clicks do I need to hit to export this file as “pre.pptx” (PowerPoint 2007–365 compatible) straight to my ~/Desktop folder?
Generated: 2025-09-10 13:23:27
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
from pptx import Presentation
from collections import Counter


def _gather_slide_texts(prs):
    """Collect visible text strings from every slide in the presentation."""
    slide_texts = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = shape.text
                if txt and txt.strip():
                    texts.append(txt.strip())
        slide_texts.append(texts)
    return slide_texts


def _compare_text_lists(texts1, texts2):
    """Return proportion (0-1) of slides whose text collections match case-insensitively."""
    if len(texts1) != len(texts2):
        return 0.0

    matches = 0
    for t1, t2 in zip(texts1, texts2):
        if Counter([s.lower() for s in t1]) == Counter([s.lower() for s in t2]):
            matches += 1
    return matches / len(texts1) if texts1 else 0.0


def verify_powerpoint_export(original_path, exported_path):
    """Verify that the deck was exported to exported_path faithfully.

    Scoring (progressive):
    • 0.3 – Exported file exists, is PPTX and loads successfully
    • 0.3 – Slide count matches the original file exactly
    • 0.4 – All slide texts match exactly (case-insensitive, order-agnostic)
      (partial credit awarded proportionally for partial text match)
    Returns a float between 0.0 and 1.0 and prints detailed diagnostics.
    """
    total_score = 0.0
    max_score = 1.0

    print(f"Verifying exported presentation at: {exported_path}")
    print(f"Against original presentation: {original_path}\n")

    # ----- Requirement 1: file existence & loadability (0.3) -----
    if not (os.path.exists(exported_path) and exported_path.lower().endswith(".pptx")):
        print("✗ Exported file not found or not a .pptx – 0 points awarded for export file")
        print(f"REWARD: {total_score}")
        return total_score

    try:
        exported_prs = Presentation(exported_path)
        slide_count_exported = len(exported_prs.slides)
        print(f"✓ Exported file loads with {slide_count_exported} slides (0.3)")
        total_score += 0.3
    except Exception as e:
        print(f"✗ Unable to load exported PPTX: {e} – 0 points")
        print(f"REWARD: {total_score}")
        return total_score

    # Load original presentation for comparison
    try:
        original_prs = Presentation(original_path)
    except Exception as e:
        print(f"✗ Cannot load original presentation for comparison: {e}")
        print(f"REWARD: {total_score}")
        return total_score  # cannot continue meaningful comparison

    # ----- Requirement 2: slide count equality (0.3) -----
    slide_count_original = len(original_prs.slides)
    if slide_count_original == slide_count_exported:
        print(f"✓ Slide count matches ({slide_count_original}) (0.3)")
        total_score += 0.3
    else:
        print(f"✗ Slide count mismatch – original {slide_count_original}, exported {slide_count_exported}")

    # ----- Requirement 3: text content comparison (up to 0.4) -----
    orig_texts = _gather_slide_texts(original_prs)
    exp_texts = _gather_slide_texts(exported_prs)
    proportion = _compare_text_lists(orig_texts, exp_texts)

    if proportion == 1.0 and orig_texts:
        print("✓ All slide texts match exactly (0.4)")
        total_score += 0.4
    elif proportion > 0:
        partial_points = 0.4 * proportion
        print(f"✓ Partial text match: {proportion*100:.1f}% identical (+{partial_points:.2f})")
        total_score += partial_points
    else:
        print("✗ No slide text match – 0 points")

    final_score = min(total_score, max_score)
    print(f"\nTotal Score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score


if __name__ == "__main__":
    ORIGINAL_FILE = "/home/user/almost_done_editing_but_the_marketing_team_needs_the_deck_in_native_powerpoint_format_in_libreoffice_golden.pptx"
    EXPORTED_FILE = os.path.expanduser("~/Desktop/pre.pptx")

    verify_powerpoint_export(ORIGINAL_FILE, EXPORTED_FILE)
