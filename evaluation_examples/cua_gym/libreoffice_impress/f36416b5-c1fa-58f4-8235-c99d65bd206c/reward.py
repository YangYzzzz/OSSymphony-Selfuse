"""
FINAL REWARD SCRIPT - SUCCESS
Task: While tidying up a massive deck in LibreOffice Impress, I noticed slide 126 doesn’t match the rest. Could you walk me through switching the title text to the palette color “Orange 4” (#FF6600) and completely disabling the shadow effect for that title?
Generated: 2025-09-10 18:09:04
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from lxml import etree

# -----------------------------------------------------------------------------
# Reward Script for LibreOffice Impress Task Verification
# Task: Ensure the title on slide 126 is colored “Orange 4” (#FF6600) and has
#       NO shadow effect applied.
# -----------------------------------------------------------------------------
# Scoring (progressive):
#   • Correct title colour  → 0.5 points
#   • Shadow effect absent → 0.5 points
# Only when BOTH conditions are met will the script return 1.0.
# -----------------------------------------------------------------------------
EXPECTED_RGB = RGBColor(255, 102, 0)          # Hex #FF6600 (Orange 4)
SLIDE_INDEX   = 125                           # 0-based index for slide 126
FILE_PATH     = "/home/user/while_tidying_up_a_massive_deck_in_libreoffice_impress_i_noticed_slide_126_doesnt_match_the_rest_cou_golden.pptx"

# --------------------------- Helper Verification ---------------------------

def verify_title_color(title_shape):
    """Verify every text run in the title uses the expected RGB colour."""
    all_correct = True
    run_found   = False
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run_found = True
            font = run.font
            # Ensure an explicit RGB value is set
            if font.color is None or font.color.type != 1:
                print(f"    ✗ Run '{run.text}' lacks an explicit RGB colour")
                all_correct = False
            else:
                rgb = font.color.rgb
                print(f"    Run '{run.text}' colour: {rgb}")
                if rgb != EXPECTED_RGB:
                    print("      ✗ Colour does not match Orange 4 (#FF6600)")
                    all_correct = False
    if not run_found:
        print("    ✗ No text runs detected in title shape")
        return False
    if all_correct:
        print("    ✓ All title runs have the correct Orange 4 colour")
    return all_correct

def verify_no_shadow(title_shape):
    """Verify that NO shadow XML elements exist within the title shape."""
    shadow_xpath = (
        './/*[contains(local-name(), "shdw") or contains(local-name(), "Shdw") or '
        'contains(local-name(), "outerShdw") or contains(local-name(), "innerShdw")]'
    )
    shadow_elems = title_shape._element.xpath(shadow_xpath)
    if shadow_elems:
        print(f"    ✗ Detected {len(shadow_elems)} shadow-related XML elements")
        for el in shadow_elems[:5]:               # Show up to 5 snippets
            snippet = etree.tostring(el, encoding='unicode').strip()
            print("      ", snippet[:120], "...")
        return False
    print("    ✓ No shadow elements found in title shape XML")
    return True

# ------------------------------- Main Check ---------------------------------

def verify_task(file_path):
    print(f"Starting verification for: {file_path}\n")
    total_score = 0.0

    # --- Preliminary: file presence & loading (no points) ---
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0
    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded successfully with {len(prs.slides)} slides\n")
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        return 0.0

    # --- Ensure slide 126 exists ---
    if len(prs.slides) <= SLIDE_INDEX:
        print("✗ Slide 126 (index 125) is missing from the presentation")
        return 0.0

    slide = prs.slides[SLIDE_INDEX]
    title_shape = slide.shapes.title
    if title_shape is None:
        print("✗ Slide 126 lacks a title placeholder")
        return 0.0
    print("✓ Title placeholder found on slide 126\n")

    # --- Requirement 1: Correct colour (0.5) ---
    if verify_title_color(title_shape):
        total_score += 0.5
        print("  → Colour requirement satisfied (+0.5)\n")
    else:
        print("  → Colour requirement NOT satisfied (+0.0)\n")

    # --- Requirement 2: No shadow (0.5) ---
    if verify_no_shadow(title_shape):
        total_score += 0.5
        print("  → Shadow requirement satisfied (+0.5)\n")
    else:
        print("  → Shadow requirement NOT satisfied (+0.0)\n")

    final_score = min(total_score, 1.0)
    print(f"Final Score: {final_score}/1.0")
    return final_score

# ---------------------------- Execute & Output -----------------------------
if __name__ == "__main__":
    reward = verify_task(FILE_PATH)
    print(f"REWARD: {reward}")
