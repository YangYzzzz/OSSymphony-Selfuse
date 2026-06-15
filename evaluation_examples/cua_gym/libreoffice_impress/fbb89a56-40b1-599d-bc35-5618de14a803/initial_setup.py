"""
Initial Setup: Create a Biology Review Game presentation with 10 slides.
Slide 8 has only the title 'Review Question 3' and is otherwise empty.
Task ID: impress_teach_035
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_035'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_text(slide, text, font_size=32, bold=True, color=None):
    """Set the title placeholder text with formatting."""
    title = slide.shapes.title
    title.text = text
    for run in title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color


def add_content_textbox(slide, left, top, width, height, text, font_size=18,
                        bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Add a text box with formatted content."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    for run in p.runs:
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    dark_blue = RGBColor(0x1B, 0x3A, 0x5C)
    white = RGBColor(0xFF, 0xFF, 0xFF)
    light_gray = RGBColor(0x44, 0x72, 0xC4)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Biology Review Game"
    slide1.placeholders[1].text = "Test Your Knowledge - Mrs. Patterson's AP Biology"

    # --- Slide 2: Instructions ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_content_textbox(slide2, Inches(1), Inches(0.5), Inches(8), Inches(1),
                        "How to Play", font_size=28, bold=True, color=dark_blue,
                        alignment=PP_ALIGN.CENTER)
    instructions = (
        "1. Each slide contains a review question from our unit on Cell Biology.\n"
        "2. Try to answer the question before revealing the hidden answer.\n"
        "3. Select the white area below each question to reveal the answer.\n"
        "4. Keep track of how many you get correct!"
    )
    add_content_textbox(slide2, Inches(1), Inches(1.8), Inches(8), Inches(4),
                        instructions, font_size=16)

    # --- Slide 3: Review Question 1 ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_content_textbox(slide3, Inches(1), Inches(0.3), Inches(8), Inches(0.8),
                        "Review Question 1", font_size=28, bold=True, color=dark_blue,
                        alignment=PP_ALIGN.CENTER)
    add_content_textbox(slide3, Inches(1), Inches(1.5), Inches(8), Inches(1.5),
                        "What is the basic structural and functional unit of all living organisms?",
                        font_size=20, alignment=PP_ALIGN.CENTER)
    # Hidden answer
    txBox = add_content_textbox(slide3, Inches(2.5), Inches(4), Inches(5), Inches(1),
                                "The Cell", font_size=20, bold=True,
                                color=white, alignment=PP_ALIGN.CENTER)
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = white

    # --- Slide 4: Review Question 2 ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_content_textbox(slide4, Inches(1), Inches(0.3), Inches(8), Inches(0.8),
                        "Review Question 2", font_size=28, bold=True, color=dark_blue,
                        alignment=PP_ALIGN.CENTER)
    add_content_textbox(slide4, Inches(1), Inches(1.5), Inches(8), Inches(1.5),
                        "Name the organelle responsible for producing ATP, the cell's main energy currency.",
                        font_size=20, alignment=PP_ALIGN.CENTER)
    txBox = add_content_textbox(slide4, Inches(2.5), Inches(4), Inches(5), Inches(1),
                                "Mitochondria", font_size=20, bold=True,
                                color=white, alignment=PP_ALIGN.CENTER)
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = white

    # --- Slide 5: Matching Activity ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_content_textbox(slide5, Inches(1), Inches(0.3), Inches(8), Inches(0.8),
                        "Organelle Matching", font_size=28, bold=True, color=dark_blue,
                        alignment=PP_ALIGN.CENTER)
    matching = (
        "Match each organelle with its function:\n\n"
        "A. Ribosome          1. Protein folding and transport\n"
        "B. Golgi Apparatus   2. Protein synthesis\n"
        "C. Endoplasmic Reticulum  3. Packaging and shipping proteins\n"
        "D. Lysosome          4. Cellular digestion"
    )
    add_content_textbox(slide5, Inches(1), Inches(1.5), Inches(8), Inches(4),
                        matching, font_size=16)

    # --- Slide 6: True or False ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_content_textbox(slide6, Inches(1), Inches(0.3), Inches(8), Inches(0.8),
                        "True or False", font_size=28, bold=True, color=dark_blue,
                        alignment=PP_ALIGN.CENTER)
    add_content_textbox(slide6, Inches(1), Inches(1.5), Inches(8), Inches(1.5),
                        "Plant cells and animal cells both contain chloroplasts.",
                        font_size=20, alignment=PP_ALIGN.CENTER)
    txBox = add_content_textbox(slide6, Inches(2.5), Inches(4), Inches(5), Inches(1),
                                "False - Only plant cells contain chloroplasts",
                                font_size=20, bold=True, color=white,
                                alignment=PP_ALIGN.CENTER)
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = white

    # --- Slide 7: Diagram Label ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_content_textbox(slide7, Inches(1), Inches(0.3), Inches(8), Inches(0.8),
                        "Cell Membrane Structure", font_size=28, bold=True, color=dark_blue,
                        alignment=PP_ALIGN.CENTER)
    add_content_textbox(slide7, Inches(1), Inches(1.5), Inches(8), Inches(3),
                        "The cell membrane is composed of a phospholipid bilayer with embedded proteins.\n\n"
                        "Key components:\n"
                        "- Phospholipid molecules (hydrophilic heads, hydrophobic tails)\n"
                        "- Integral proteins (span the entire membrane)\n"
                        "- Peripheral proteins (attached to surface)\n"
                        "- Cholesterol molecules (regulate fluidity)",
                        font_size=16)

    # --- Slide 8: Review Question 3 --- (EMPTY except title)
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    add_content_textbox(slide8, Inches(1), Inches(0.3), Inches(8), Inches(0.8),
                        "Review Question 3", font_size=28, bold=True, color=dark_blue,
                        alignment=PP_ALIGN.CENTER)
    # Intentionally empty - the agent must add the question and hidden answer

    # --- Slide 9: Bonus Question ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    add_content_textbox(slide9, Inches(1), Inches(0.3), Inches(8), Inches(0.8),
                        "Bonus Question", font_size=28, bold=True, color=dark_blue,
                        alignment=PP_ALIGN.CENTER)
    add_content_textbox(slide9, Inches(1), Inches(1.5), Inches(8), Inches(1.5),
                        "What scientist first observed cells under a microscope in 1665?",
                        font_size=20, alignment=PP_ALIGN.CENTER)
    txBox = add_content_textbox(slide9, Inches(2.5), Inches(4), Inches(5), Inches(1),
                                "Robert Hooke", font_size=20, bold=True,
                                color=white, alignment=PP_ALIGN.CENTER)
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = white

    # --- Slide 10: Score Tracker ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    add_content_textbox(slide10, Inches(1), Inches(0.3), Inches(8), Inches(0.8),
                        "Your Score", font_size=28, bold=True, color=dark_blue,
                        alignment=PP_ALIGN.CENTER)
    add_content_textbox(slide10, Inches(1), Inches(1.8), Inches(8), Inches(3),
                        "How did you do?\n\n"
                        "7-8 Correct: Excellent! You're ready for the exam.\n"
                        "5-6 Correct: Good job! Review a few more topics.\n"
                        "3-4 Correct: Keep studying - focus on weak areas.\n"
                        "0-2 Correct: Schedule a review session with Mrs. Patterson.",
                        font_size=18, alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
