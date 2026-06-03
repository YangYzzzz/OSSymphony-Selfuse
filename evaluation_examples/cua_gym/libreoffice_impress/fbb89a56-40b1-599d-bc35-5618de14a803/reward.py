"""
Reward Script: Fill-in-the-blank review slide on slide 8
Task ID: impress_teach_035
Domain: libreoffice_impress
Scoring:
  Component 1 (0.35): Question text box with blank on slide 8
  Component 2 (0.25): Answer text box with "Photosynthesis"
  Component 3 (0.25): Answer text color is white (#FFFFFF)
  Component 4 (0.15): Answer text box background is white (#FFFFFF) solid fill
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_035'


def get_all_text_shapes(slide):
    """Extract all shapes with text frames, including nested groups."""
    def extract(shape):
        results = []
        if hasattr(shape, "text") and hasattr(shape, "text_frame"):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.extend(extract(sub))
        return results
    out = []
    for shape in slide.shapes:
        out.extend(extract(shape))
    return out


def persist_app_state(domain):
    """Try to save any unsaved GUI state via Ctrl+S."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 8 slides
    if len(prs.slides) < 8:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 8")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[7]  # 0-indexed, slide 8
    text_shapes = get_all_text_shapes(slide)

    # Identify question shape and answer shape by content
    question_shape = None
    answer_shape = None

    for shape in text_shapes:
        full_text = shape.text_frame.text.strip() if shape.has_text_frame else ""
        lower_text = full_text.lower()

        # Question: contains the blank line indicator and "plants" or "sunlight" or "energy"
        if "_____" in full_text and ("plants" in lower_text or "sunlight" in lower_text or "energy" in lower_text):
            question_shape = shape

        # Answer: contains "photosynthesis" (case-insensitive)
        if "photosynthesis" in lower_text and "_____" not in full_text:
            answer_shape = shape

    # Component 1: Question text box with fill-in-the-blank (0.35 points)
    try:
        if question_shape is not None:
            q_text = question_shape.text_frame.text.strip()
            if "_____" in q_text:
                print(f"PASS: Component 1 -- Question text found: '{q_text[:60]}...' (0.35 pts)")
                total_score += 0.35
            else:
                print(f"FAIL: Component 1 -- Question shape found but no blank underscores")
        else:
            print(f"FAIL: Component 1 -- No question text box with blank found on slide 8")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Answer text box contains "Photosynthesis" (0.25 points)
    try:
        if answer_shape is not None:
            a_text = answer_shape.text_frame.text.strip()
            if "photosynthesis" in a_text.lower():
                print(f"PASS: Component 2 -- Answer text found: '{a_text}' (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 2 -- Answer shape found but text is '{a_text}'")
        else:
            print(f"FAIL: Component 2 -- No answer text box with 'Photosynthesis' found on slide 8")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Answer text is white (#FFFFFF) (0.25 points)
    try:
        if answer_shape is not None:
            all_runs_white = True
            any_run_found = False
            for para in answer_shape.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        any_run_found = True
                        try:
                            if run.font.color.type is not None:
                                rgb = str(run.font.color.rgb).upper()
                                if rgb != "FFFFFF":
                                    all_runs_white = False
                                    print(f"FAIL: Component 3 -- Run text '{run.text}' has color {rgb}, expected FFFFFF")
                            else:
                                all_runs_white = False
                                print(f"FAIL: Component 3 -- Run text '{run.text}' has no explicit color (inherited/theme)")
                        except Exception:
                            all_runs_white = False
                            print(f"FAIL: Component 3 -- Could not read color for run '{run.text}'")

            if any_run_found and all_runs_white:
                print(f"PASS: Component 3 -- Answer text color is white #FFFFFF (0.25 pts)")
                total_score += 0.25
            elif not any_run_found:
                print(f"FAIL: Component 3 -- No non-empty runs found in answer shape")
        else:
            print(f"FAIL: Component 3 -- No answer shape found, cannot check text color")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Answer text box has white (#FFFFFF) solid fill background (0.15 points)
    try:
        if answer_shape is not None:
            fill = answer_shape.fill
            if fill.type == 1:  # SOLID fill
                try:
                    bg_rgb = str(fill.fore_color.rgb).upper()
                    if bg_rgb == "FFFFFF":
                        print(f"PASS: Component 4 -- Answer box background is white #FFFFFF solid fill (0.15 pts)")
                        total_score += 0.15
                    else:
                        print(f"FAIL: Component 4 -- Answer box background color is {bg_rgb}, expected FFFFFF")
                except Exception as e2:
                    print(f"FAIL: Component 4 -- Could not read fill color: {e2}")
            else:
                print(f"FAIL: Component 4 -- Answer box fill type is {fill.type}, expected SOLID (1)")
        else:
            print(f"FAIL: Component 4 -- No answer shape found, cannot check background")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
