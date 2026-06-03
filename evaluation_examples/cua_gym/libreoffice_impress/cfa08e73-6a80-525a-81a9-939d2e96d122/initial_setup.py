"""
Initial Setup: Create a 5-slide presentation with a bar chart on slide 4
that has mixed fonts and sizes across chart text elements.
Task ID: impress_tct_065
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_065'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_font(font_obj, name, size_pt, bold=False, italic=False, color_rgb=None):
    """Helper to set font properties."""
    font_obj.name = name
    font_obj.size = Pt(size_pt)
    font_obj.bold = bold
    font_obj.italic = italic
    if color_rgb:
        font_obj.color.rgb = RGBColor(*color_rgb)


def add_text_slide(prs, layout_idx, title_text, body_lines):
    """Add a slide with title and bullet body text."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title_text
    if body_lines and len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, line in enumerate(body_lines):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                p = tf.add_paragraph()
                p.text = line
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Consistent Charts Initiative"
    slide1.placeholders[1].text = "Q4 2025 Data Visualization Standards Review"

    # --- Slide 2: Overview ---
    slide2 = add_text_slide(prs, 1, "Project Overview", [
        "Standardize chart formatting across all departments",
        "Ensure readability and accessibility compliance",
        "Align with corporate branding guidelines",
        "Reduce manual formatting effort by 60%",
        "Target completion: End of Q1 2026",
    ])

    # --- Slide 3: Key Metrics ---
    slide3 = add_text_slide(prs, 1, "Key Performance Metrics", [
        "Revenue grew 23% year-over-year to $4.2M",
        "Customer acquisition cost decreased by 15%",
        "Net promoter score improved from 42 to 58",
        "Employee satisfaction at 87% (up from 79%)",
        "Market share expanded to 18.5% in core segment",
    ])

    # --- Slide 4: Bar Chart with MIXED fonts/sizes (initial state) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Add a title textbox for the slide
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10), Inches(0.6))
    tf = txBox.text_frame
    tf.paragraphs[0].text = "Regional Sales Performance"
    for run in tf.paragraphs[0].runs:
        run.font.size = Pt(28)
        run.font.bold = True

    # Create bar chart data
    chart_data = CategoryChartData()
    chart_data.categories = ['North', 'South', 'East', 'West', 'Central']
    chart_data.add_series('Q3 2025', (342, 278, 456, 189, 315))
    chart_data.add_series('Q4 2025', (398, 312, 489, 225, 367))

    chart_frame = slide4.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8), Inches(1.0), Inches(11.0), Inches(5.8),
        chart_data,
    )
    chart = chart_frame.chart

    # Chart title - set to 14pt Arial (NOT 10pt Calibri)
    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = "Quarterly Sales by Region"
    # Set on paragraph default AND on run
    title_font = chart.chart_title.text_frame.paragraphs[0].font
    set_font(title_font, 'Arial', 14, bold=True)
    for run in chart.chart_title.text_frame.paragraphs[0].runs:
        set_font(run.font, 'Arial', 14, bold=True)

    # Value axis title - 11pt Times New Roman
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.paragraphs[0].text = "Revenue (thousands USD)"
    va_title_font = chart.value_axis.axis_title.text_frame.paragraphs[0].font
    set_font(va_title_font, 'Times New Roman', 11)
    for run in chart.value_axis.axis_title.text_frame.paragraphs[0].runs:
        set_font(run.font, 'Times New Roman', 11)

    # Category axis title - 11pt Times New Roman
    chart.category_axis.has_title = True
    chart.category_axis.axis_title.text_frame.paragraphs[0].text = "Region"
    ca_title_font = chart.category_axis.axis_title.text_frame.paragraphs[0].font
    set_font(ca_title_font, 'Times New Roman', 11)
    for run in chart.category_axis.axis_title.text_frame.paragraphs[0].runs:
        set_font(run.font, 'Times New Roman', 11)

    # Value axis tick labels - 8pt Verdana
    chart.value_axis.tick_labels.font.name = 'Verdana'
    chart.value_axis.tick_labels.font.size = Pt(8)

    # Category axis tick labels - 8pt Verdana
    chart.category_axis.tick_labels.font.name = 'Verdana'
    chart.category_axis.tick_labels.font.size = Pt(8)

    # Data labels - 8pt Courier New
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_value = True
    plot.data_labels.font.name = 'Courier New'
    plot.data_labels.font.size = Pt(8)

    # Legend - 11pt Arial
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.name = 'Arial'
    chart.legend.font.size = Pt(11)

    # --- Slide 5: Next Steps ---
    slide5 = add_text_slide(prs, 1, "Next Steps & Timeline", [
        "Phase 1: Audit all existing charts (Jan 15 - Jan 31)",
        "Phase 2: Define standard templates (Feb 1 - Feb 15)",
        "Phase 3: Apply formatting across departments (Feb 16 - Mar 15)",
        "Phase 4: Training sessions for team leads (Mar 16 - Mar 31)",
        "Final review and sign-off by April 10, 2026",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
