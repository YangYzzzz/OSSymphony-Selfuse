"""
Initial Setup: Over-animated presentation with entrance + emphasis animations on slide 7
Task ID: impress_fix_062
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import zipfile
import shutil
import copy
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_062'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_paragraph(tf, text, level=0, font_size=Pt(18), bold=False, color=None):
    """Add a paragraph to a text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.space_after = Pt(6)
    for run in p.runs:
        run.font.size = font_size
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
    return p


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q4 2025 Strategic Growth Plan"
    slide1.placeholders[1].text = "Prepared by the Executive Strategy Team\nNovember 2025"
    for run in slide1.shapes.title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # --- Slide 2: Market Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Market Overview"
    tf2 = slide2.placeholders[1].text_frame
    tf2.clear()
    tf2.paragraphs[0].text = "Global SaaS market projected to reach $908B by 2030"
    for item in [
        "Enterprise segment growing at 14.2% CAGR",
        "Mid-market adoption accelerating in APAC region",
        "Competitive landscape shifting toward AI-native platforms",
        "Customer retention rates improving industry-wide to 92%",
    ]:
        add_paragraph(tf2, item, level=0, font_size=Pt(18))

    # --- Slide 3: Revenue Performance ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Revenue Performance YTD"
    tf3 = slide3.placeholders[1].text_frame
    tf3.clear()
    tf3.paragraphs[0].text = "Total Revenue: $47.3M (up 23% YoY)"
    for item in [
        "Subscription Revenue: $38.1M (+28% YoY)",
        "Professional Services: $6.8M (+12% YoY)",
        "Hardware & Licensing: $2.4M (-5% YoY)",
        "Net Revenue Retention: 118%",
        "New Logo Acquisition: 142 accounts",
    ]:
        add_paragraph(tf3, item, level=0, font_size=Pt(18))

    # --- Slide 4: Customer Segments ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Customer Segment Analysis"
    tf4 = slide4.placeholders[1].text_frame
    tf4.clear()
    tf4.paragraphs[0].text = "Enterprise (>1000 employees): 34% of revenue"
    for item in [
        "Mid-Market (100-999 employees): 41% of revenue",
        "SMB (<100 employees): 25% of revenue",
        "Fastest growing: Mid-Market APAC (+38% QoQ)",
        "Highest retention: Enterprise NA (97.2%)",
    ]:
        add_paragraph(tf4, item, level=0, font_size=Pt(18))

    # --- Slide 5: Product Roadmap ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Product Roadmap Highlights"
    tf5 = slide5.placeholders[1].text_frame
    tf5.clear()
    tf5.paragraphs[0].text = "AI Assistant launch: January 2026"
    for item in [
        "Advanced Analytics Dashboard: February 2026",
        "Mobile App v3.0 with offline mode: March 2026",
        "Enterprise SSO & SCIM provisioning: Q1 2026",
        "Marketplace integrations (50+ partners): Q2 2026",
    ]:
        add_paragraph(tf5, item, level=0, font_size=Pt(18))

    # --- Slide 6: Team Growth ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Team & Organizational Growth"
    tf6 = slide6.placeholders[1].text_frame
    tf6.clear()
    tf6.paragraphs[0].text = "Current headcount: 312 employees across 8 offices"
    for item in [
        "Engineering: 148 (hiring 35 more in Q1)",
        "Sales & Marketing: 89 (expanding EMEA team)",
        "Customer Success: 42 (new VP hired from Salesforce)",
        "G&A: 33 (streamlining with automation)",
    ]:
        add_paragraph(tf6, item, level=0, font_size=Pt(18))

    # --- Slide 7: Key Initiatives (the animated slide) ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Key Strategic Initiatives for Q4"
    tf7 = slide7.placeholders[1].text_frame
    tf7.clear()
    bullet_texts = [
        "Launch AI-powered customer onboarding to reduce time-to-value by 40%",
        "Expand into the Japanese market with localized product and support team",
        "Achieve SOC 2 Type II and ISO 27001 certifications before year end",
        "Migrate remaining legacy infrastructure to Kubernetes by December",
        "Establish strategic partnership program with top 10 system integrators",
    ]
    tf7.paragraphs[0].text = bullet_texts[0]
    for txt in bullet_texts[1:]:
        add_paragraph(tf7, txt, level=0, font_size=Pt(20))

    # Save the base presentation first
    prs.save(OUTPUT)
    print(f'Base presentation created: {OUTPUT}')

    # Now inject animations into slide 7 via XML manipulation
    inject_animations(OUTPUT)
    print(f'Animations injected into slide 7')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


def inject_animations(pptx_path):
    """Inject Fly In entrance + Pulse emphasis animations for each bullet on slide 7."""
    tmp_path = pptx_path + '.tmp'
    shutil.copy(pptx_path, tmp_path)

    # Namespaces used in OOXML presentation XML
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }

    with zipfile.ZipFile(tmp_path, 'r') as zin:
        with zipfile.ZipFile(pptx_path, 'w') as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == 'ppt/slides/slide7.xml':
                    data = add_animations_to_slide(data, nsmap)
                zout.writestr(item, data)

    os.remove(tmp_path)


def add_animations_to_slide(slide_xml_bytes, nsmap):
    """Add Fly In (entrance) + Pulse (emphasis) animations for 5 bullet paragraphs."""
    root = etree.fromstring(slide_xml_bytes)

    # Find the content placeholder shape (the one with bullet text, not the title)
    p_ns = nsmap['p']
    a_ns = nsmap['a']

    # Find shape elements
    spTree = root.find(f'.//{{{p_ns}}}cSld/{{{p_ns}}}spTree')
    if spTree is None:
        print("WARNING: Could not find spTree")
        return slide_xml_bytes

    # Find content placeholder (type "body" or idx=1)
    content_sp = None
    for sp in spTree.findall(f'{{{p_ns}}}sp'):
        nvSpPr = sp.find(f'{{{p_ns}}}nvSpPr')
        if nvSpPr is not None:
            nvPr = nvSpPr.find(f'{{{p_ns}}}nvPr')
            if nvPr is not None:
                ph = nvPr.find(f'{{{p_ns}}}ph')
                if ph is not None:
                    ph_type = ph.get('type', '')
                    ph_idx = ph.get('idx', '')
                    if ph_type == 'body' or ph_idx == '1':
                        content_sp = sp
                        break

    if content_sp is None:
        print("WARNING: Could not find content placeholder")
        return slide_xml_bytes

    # Get the shape ID for the content placeholder
    cNvPr = content_sp.find(f'{{{p_ns}}}nvSpPr/{{{p_ns}}}cNvPr')
    shape_id = cNvPr.get('id') if cNvPr is not None else '3'

    # Build the timing/animation XML
    # We need 5 entrance (fly in) + 5 emphasis (pulse) animations
    # Each bullet paragraph is targeted by (shapeId, paragraph index)
    timing_xml = build_timing_xml(shape_id, nsmap)

    # Remove existing timing if any
    existing_timing = root.find(f'{{{p_ns}}}timing')
    if existing_timing is not None:
        root.remove(existing_timing)

    # Add timing element
    timing_elem = etree.fromstring(timing_xml)
    root.append(timing_elem)

    return etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)


def build_timing_xml(shape_id, nsmap):
    """Build OOXML timing XML for 5 Fly In entrance + 5 Pulse emphasis animations."""
    p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # Each bullet gets:
    #   1. Fly In entrance animation (anim:fly, dir="b" for from bottom)
    #   2. Pulse emphasis animation (anim:pulse / animScale)
    # Animations are sequenced: click -> fly in bullet 1, then pulse bullet 1,
    #   click -> fly in bullet 2, then pulse bullet 2, etc.

    # Build child sequence nodes for each bullet
    child_sequences = []
    for i in range(5):
        para_idx = i  # 0-based paragraph index

        # Node ID tracking (each anim node needs unique IDs)
        # We'll use a simple scheme: bullet i uses base id 2 + i*20
        base_id = 2 + i * 20

        # Determine if first click or subsequent
        if i == 0:
            # First sequence starts on click
            seq_start = f'''
            <p:stCondLst>
              <p:cond evt="onClick" delay="0">
                <p:tgtEl><p:sldTgt/></p:tgtEl>
              </p:cond>
            </p:stCondLst>'''
        else:
            seq_start = f'''
            <p:stCondLst>
              <p:cond evt="onClick" delay="0">
                <p:tgtEl><p:sldTgt/></p:tgtEl>
              </p:cond>
            </p:stCondLst>'''

        # Fly In entrance animation for this bullet
        fly_in_node = f'''
        <p:par>
          <p:cTn id="{base_id + 1}" presetID="2" presetClass="entr" presetSubtype="4"
                 fill="hold" nodeType="clickEffect">
            <p:stCondLst><p:cond delay="0"/></p:stCondLst>
            <p:childTnLst>
              <p:set>
                <p:cBhvr>
                  <p:cTn id="{base_id + 2}" dur="1" fill="hold">
                    <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                  </p:cTn>
                  <p:tgtEl>
                    <p:spTgt spid="{shape_id}">
                      <p:txEl><p:pRg st="{para_idx}" end="{para_idx}"/></p:txEl>
                    </p:spTgt>
                  </p:tgtEl>
                  <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                </p:cBhvr>
                <p:to><p:strVal val="visible"/></p:to>
              </p:set>
              <p:anim calcmode="lin" valueType="num">
                <p:cBhvr additive="base">
                  <p:cTn id="{base_id + 3}" dur="500" fill="hold"/>
                  <p:tgtEl>
                    <p:spTgt spid="{shape_id}">
                      <p:txEl><p:pRg st="{para_idx}" end="{para_idx}"/></p:txEl>
                    </p:spTgt>
                  </p:tgtEl>
                  <p:attrNameLst><p:attrName>ppt_y</p:attrName></p:attrNameLst>
                </p:cBhvr>
                <p:tavLst>
                  <p:tav tm="0"><p:val><p:strVal val="1+#ppt_h/2"/></p:val></p:tav>
                  <p:tav tm="100000"><p:val><p:strVal val="#ppt_y"/></p:val></p:tav>
                </p:tavLst>
              </p:anim>
            </p:childTnLst>
          </p:cTn>
        </p:par>'''

        # Pulse emphasis animation for this bullet
        pulse_node = f'''
        <p:par>
          <p:cTn id="{base_id + 5}" presetID="10" presetClass="emph" presetSubtype="0"
                 fill="hold" nodeType="afterEffect">
            <p:stCondLst><p:cond delay="0"/></p:stCondLst>
            <p:childTnLst>
              <p:animScale>
                <p:cBhvr>
                  <p:cTn id="{base_id + 6}" dur="500" autoRev="1" fill="hold"/>
                  <p:tgtEl>
                    <p:spTgt spid="{shape_id}">
                      <p:txEl><p:pRg st="{para_idx}" end="{para_idx}"/></p:txEl>
                    </p:spTgt>
                  </p:tgtEl>
                </p:cBhvr>
                <p:by x="110000" y="110000"/>
              </p:animScale>
            </p:childTnLst>
          </p:cTn>
        </p:par>'''

        # Each click group: entrance then emphasis after
        click_group = f'''
      <p:seq concurrent="1" nextAc="seek">
        <p:cTn id="{base_id}" fill="hold" nodeType="clickSeq">
          {seq_start}
          <p:endSync evt="end" delay="0">
            <p:rtn val="all"/>
          </p:endSync>
          <p:childTnLst>
            <p:par>
              <p:cTn id="{base_id + 10}" fill="hold">
                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                <p:childTnLst>
                  <p:par>
                    <p:cTn id="{base_id + 11}" fill="hold">
                      <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                      <p:childTnLst>
                        {fly_in_node}
                        {pulse_node}
                      </p:childTnLst>
                    </p:cTn>
                  </p:par>
                </p:childTnLst>
              </p:cTn>
            </p:par>
          </p:childTnLst>
        </p:cTn>
        <p:prevCondLst>
          <p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
        </p:prevCondLst>
        <p:nextCondLst>
          <p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
        </p:nextCondLst>
      </p:seq>'''

        child_sequences.append(click_group)

    all_sequences = '\n'.join(child_sequences)

    timing_xml = f'''<p:timing xmlns:p="{p_ns}" xmlns:a="{a_ns}">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          {all_sequences}
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
  <p:bldLst>
    <p:bldP spid="{shape_id}" grpId="0" build="p"/>
  </p:bldLst>
</p:timing>'''

    return timing_xml


create_presentation()
