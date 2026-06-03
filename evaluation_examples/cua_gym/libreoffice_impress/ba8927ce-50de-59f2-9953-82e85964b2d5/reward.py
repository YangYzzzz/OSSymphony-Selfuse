"""
FINAL REWARD SCRIPT - SUCCESS
Task: I’m cleaning up a massive training deck in LibreOffice Impress (it’s 300 slides long!). On slide 274 there’s the word "Website" in the closing bullet list—right now it’s just plain text. I need that exact word to link straight to https://example.com/docs. How do I set up that hyperlink?
Generated: 2025-09-10 21:09:49
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation

def verify_impress_hyperlink(file_path: str) -> float:
    """Verify that on slide 274 (index 273) the exact word 'Website'
    is present and hyperlinks to https://example.com/docs.

    Progressive scoring:
        - 0.4 pts for finding the word 'Website' on the correct slide
        - 0.6 pts for the hyperlink being exactly the expected URL
    Returns a score between 0.0 and 1.0 and prints detailed feedback.
    """

    expected_url = "https://example.com/docs"
    total_score = 0.0
    max_score = 1.0

    print(f"Verifying presentation: {file_path}\n")

    # --- 1. File existence and loading (no points, prerequisite) ---
    if not os.path.exists(file_path):
        print("✗ File not found.")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Failed to load PPTX: {e}")
        print("REWARD: 0.0")
        return 0.0

    # --- 2. Verify slide count is sufficient (no direct points) ---
    target_index = 273  # zero-based index for slide 274
    if len(prs.slides) <= target_index:
        print(f"✗ Presentation has only {len(prs.slides)} slides; needs ≥ 274.")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[target_index]

    # --- 3. Search for the word 'Website' and its hyperlink ---
    website_text_found = False
    correct_hyperlink = False

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text.strip().lower() == "website":
                    website_text_found = True
                    url = run.hyperlink.address
                    print(f"Found 'Website' run with hyperlink: {url}")
                    if url and url.strip() == expected_url:
                        correct_hyperlink = True

    # --- 4. Scoring based on actual verification ---
    if website_text_found:
        print("✓ The word 'Website' was found on slide 274 (+0.4)")
        total_score += 0.4
        if correct_hyperlink:
            print("✓ Hyperlink correctly set to expected URL (+0.6)")
            total_score += 0.6
        else:
            print("✗ Hyperlink is missing or incorrect (+0.0)")
    else:
        print("✗ The word 'Website' was not found on slide 274 (+0.0)")

    # Ensure the score does not exceed 1.0
    final_score = min(total_score, max_score)

    print(f"\nTotal score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score

# --------------------
# Execute verification
# --------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/im_cleaning_up_a_massive_training_deck_in_libreoffice_impress_its_300_slides_long_on_slide_274_there_golden.pptx"
    verify_impress_hyperlink(FILE_PATH)
