"""
Initial Setup: Set up a presentation with 5 slides. Slide 4 has a large rectangle
with a solid blue fill. A headshot.jpg file exists on the desktop.
Task ID: impress_design_071
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

WORKDIR = '/home/user'
TASK_ID = 'impress_design_071'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
HEADSHOT_PATH = f'{WORKDIR}/Desktop/headshot.jpg'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_headshot():
    """Create a realistic-looking headshot.jpg on the Desktop."""
    os.makedirs(f'{WORKDIR}/Desktop', exist_ok=True)
    # Create a simple portrait-style image with some color variation
    img = Image.new('RGB', (400, 500), (200, 170, 140))
    # Add some variation to make it look like a photo
    pixels = img.load()
    for y in range(500):
        for x in range(400):
            # Background gradient (top portion)
            if y < 100:
                pixels[x, y] = (80, 100, 130)
            # Simple face-like oval region
            elif 100 <= y < 400:
                cx, cy = 200, 250
                dx, dy = abs(x - cx), abs(y - cy)
                if (dx * dx) / (120 * 120) + (dy * dy) / (150 * 150) < 1:
                    # Skin tone with slight variation
                    r = min(255, 210 + (x % 10) - 5)
                    g = min(255, 180 + (y % 10) - 5)
                    b = min(255, 150 + ((x + y) % 10) - 5)
                    pixels[x, y] = (r, g, b)
                else:
                    pixels[x, y] = (80, 100, 130)
            else:
                # Shirt/clothing area
                pixels[x, y] = (40, 60, 100)
    img.save(HEADSHOT_PATH, 'JPEG', quality=85)
    print(f'Created headshot image: {HEADSHOT_PATH}')


def create_initial():
    prs = Presentation()
    # Standard 16:9 slide dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Texture & Material Design Workshop"
    slide1.placeholders[1].text = "Visual Patterns in Digital Media\nPresented by Elena Vasquez"

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Workshop Objectives"
    body2 = slide2.placeholders[1]
    tf2 = body2.text_frame
    tf2.clear()
    objectives = [
        "Understand the fundamentals of texture mapping in presentations",
        "Learn to apply picture fills with tiling patterns",
        "Explore scaling techniques for seamless visual effects",
        "Practice creating professional design mockups",
        "Review best practices for image-based backgrounds"
    ]
    for i, obj in enumerate(objectives):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = obj
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(18)

    # --- Slide 3: Theory ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Texture Mapping Concepts"
    body3 = slide3.placeholders[1]
    tf3 = body3.text_frame
    tf3.clear()
    concepts = [
        "Tiling: Repeating a pattern across a surface area",
        "Scale Factor: Controls the size of each tile repetition",
        "Offset: Adjusts starting position of the tile grid",
        "Mirroring: Alternates tile orientation for seamless edges",
        "Aspect Ratio: Maintaining proportions during scaling"
    ]
    for i, concept in enumerate(concepts):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p.text = concept
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(16)

    # --- Slide 4: Demo Slide with Blue Rectangle ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add a title text box
    title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "Demo: Apply Texture Fill Here"
    p_title.alignment = PP_ALIGN.LEFT
    for run in p_title.runs:
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x2B, 0x2B, 0x2B)

    # Large rectangle with solid blue fill
    rect = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), Inches(1.5),
        Inches(10), Inches(5)
    )
    fill = rect.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x00, 0x70, 0xC0)  # Blue fill

    # Remove outline for clean look
    rect.line.fill.background()

    # --- Slide 5: Summary ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Key Takeaways"
    body5 = slide5.placeholders[1]
    tf5 = body5.text_frame
    tf5.clear()
    takeaways = [
        "Picture fills add depth and visual interest to shapes",
        "Tile mode creates repeating patterns from a single image",
        "50% scaling doubles the number of visible tile repetitions",
        "Consistent texture application improves presentation cohesion",
        "Always test texture fills at different display resolutions"
    ]
    for i, ta in enumerate(takeaways):
        p = tf5.paragraphs[0] if i == 0 else tf5.add_paragraph()
        p.text = ta
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(18)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


# Create headshot first, then presentation
create_headshot()
create_initial()
