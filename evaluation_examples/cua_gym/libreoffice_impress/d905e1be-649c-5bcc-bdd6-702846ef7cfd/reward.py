"""
FINAL REWARD SCRIPT - SUCCESS
Task: In LibreOffice Impress, I’m working on slide 264 and need the word “Sources” to open the file located at ~/Desktop/refs.pdf when it’s clicked during the slideshow. How do I set up that hyperlink?
Generated: 2025-09-10 19:41:51
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
import traceback
from pptx import Presentation

"""
Reward Script for LibreOffice Impress Hyperlink Verification
-----------------------------------------------------------
Task: Verify that on slide 264 (1-based indexing) of the supplied presentation
file, the word "Sources" is present AND is hyper-linked to the local file
"~/Desktop/refs.pdf".  The script awards progressive points and returns a final
score between 0.0 and 1.0 (float).  A perfect implementation returns 1.0.

Scoring rubric (adds up to 1.0):
  • 0.2 – The presentation contains at least 264 slides (i.e. slide 264 exists)
  • 0.3 – Slide 264 contains text that includes the word "Sources"
  • 0.5 – A hyperlink pointing to "~/Desktop/refs.pdf" (or any case-insensitive
           variant, with or without a file:// prefix) is attached either to the
           text run that contains "Sources", the containing shape, or is present
           in the slide’s underlying relationship list.

The script performs three independent REAL checks (no hard-coding) and
calculates the score accordingly.  It prints detailed diagnostics and finally
outputs the reward in the required format:  "REWARD: X.X".
"""

# ---------------------------------------------------------------------------
# Configuration constants
# ---------------------------------------------------------------------------
EXPECTED_PATH_FRAGMENT = "desktop/refs.pdf"  # case-insensitive fragment to look for
SLIDE_INDEX = 263  # zero-based index for slide 264

# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------

def normalize_address(addr: str) -> str:
    """Lower-case, strip and drop common URI prefixes so paths are comparable."""
    if not addr:
        return ""
    addr = addr.strip().lower()
    if addr.startswith("file:///"):
        addr = addr[8:]  # drop "file:///"
    elif addr.startswith("file://"):
        addr = addr[7:]
    return addr

def addr_matches_expected(addr: str) -> bool:
    """Return True if hyperlink address points to the expected PDF file."""
    return EXPECTED_PATH_FRAGMENT in normalize_address(addr)

# ---------------------------------------------------------------------------
# Core verification routine
# ---------------------------------------------------------------------------

def verify_task(file_path: str) -> float:
    print(f"Starting verification for: {file_path}\n")

    score = 0.0
    MAX_SCORE = 1.0

    try:
        # ------- 0.  File existence & load (prerequisite – NO points) -------
        if not os.path.exists(file_path):
            print(f"✗ File not found: {file_path}")
            return 0.0

        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        print(f"✓ Presentation loaded with {slide_count} slides")

        # ------- 1.  Slide 264 existence (0.2) -------
        if SLIDE_INDEX < slide_count:
            score += 0.2
            print("✓ Slide 264 exists (0.2 points)")
            slide = prs.slides[SLIDE_INDEX]
        else:
            print("✗ Slide 264 does not exist – cannot continue")
            return score

        # ------- 2.  Look for text containing "Sources" (0.3) -------
        sources_found = False
        hyperlink_ok = False

        for shape in slide.shapes:
            # Gather visible text (if any)
            text = ""
            if hasattr(shape, "text"):
                text = shape.text
            elif hasattr(shape, "text_frame") and shape.text_frame is not None:
                text = shape.text_frame.text

            if text and "sources" in text.lower():
                sources_found = True

                # Check each run for a hyperlink
                if hasattr(shape, "text_frame") and shape.text_frame is not None:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if "sources" in run.text.lower():
                                addr = run.hyperlink.address if run.hyperlink else ""
                                if addr_matches_expected(addr):
                                    hyperlink_ok = True
                                    print(f"✓ Hyperlink found on run: {addr}")

                # Also inspect shape-level hyperlink (in case whole box is linked)
                try:
                    shape_addr = shape.hyperlink.address if shape.hyperlink else ""
                except AttributeError:
                    shape_addr = ""
                if addr_matches_expected(shape_addr):
                    hyperlink_ok = True
                    print(f"✓ Hyperlink found on shape: {shape_addr}")

        # Extra safety: inspect slide relationships for hidden hyperlinks
        for rel in slide.part.rels.values():
            try:
                if (
                    rel.reltype
                    == "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink"
                ):
                    target = str(rel.target_ref)
                    if addr_matches_expected(target):
                        hyperlink_ok = True
                        print(f"✓ Hyperlink found in relationships: {target}")
            except Exception:
                pass

        # Scoring for text presence
        if sources_found:
            score += 0.3
            print("✓ 'Sources' text found on slide (0.3 points)")
        else:
            print("✗ 'Sources' text NOT found on slide")

        # Scoring for correct hyperlink
        if hyperlink_ok:
            score += 0.5
            print("✓ Correct hyperlink to refs.pdf present (0.5 points)")
        else:
            print("✗ Hyperlink missing or incorrect destination")

        final_score = min(score, MAX_SCORE)
        print(f"\nTotal score: {final_score}/{MAX_SCORE}")
        return final_score

    except Exception as e:
        print("✗ Exception during verification:", e)
        traceback.print_exc()
        return 0.0

# ---------------------------------------------------------------------------
# Entry point – adjust the path if necessary
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    # Path to the file under evaluation (provided by the task environment)
    FILE_PATH = "/home/user/in_libreoffice_impress_im_working_on_slide_264_and_need_the_word_sources_to_open_the_file_located_at_golden.pptx"

    reward = verify_task(FILE_PATH)
    print(f"REWARD: {reward}")
