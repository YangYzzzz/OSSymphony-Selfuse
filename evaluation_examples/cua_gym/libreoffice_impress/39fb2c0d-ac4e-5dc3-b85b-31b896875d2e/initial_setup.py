"""
Initial Setup: Build a slide reveal animation presentation
Task ID: impress_gf2_044
Domain: libreoffice_impress

Creates a 10-slide 'Story_Presentation.pptx' with slide 8 containing:
- A background illustration image (full-slide)
- A white rectangle covering entire slide area (on top)
- A narration text box
No animations applied.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_044'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
BG_IMAGE = f'{WORKDIR}/_story_illustration.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_illustration():
    """Create a colorful illustration image for the story slide."""
    from PIL import Image, ImageDraw, ImageFont

    width, height = 1280, 720
    img = Image.new('RGB', (width, height), (30, 60, 90))
    draw = ImageDraw.Draw(img)

    # Draw a landscape scene
    # Sky gradient (already blue-ish background)
    # Sun
    draw.ellipse([900, 50, 1020, 170], fill=(255, 200, 50))
    # Sun rays
    for angle_offset in range(0, 360, 30):
        import math
        cx, cy = 960, 110
        r1, r2 = 70, 100
        rad = math.radians(angle_offset)
        x1 = cx + int(r1 * math.cos(rad))
        y1 = cy + int(r1 * math.sin(rad))
        x2 = cx + int(r2 * math.cos(rad))
        y2 = cy + int(r2 * math.sin(rad))
        draw.line([(x1, y1), (x2, y2)], fill=(255, 220, 80), width=3)

    # Mountains
    draw.polygon([(0, 500), (200, 250), (400, 500)], fill=(60, 100, 60))
    draw.polygon([(250, 500), (500, 200), (750, 500)], fill=(50, 90, 50))
    draw.polygon([(600, 500), (850, 280), (1100, 500)], fill=(70, 110, 70))
    draw.polygon([(900, 500), (1100, 300), (1280, 500)], fill=(55, 95, 55))

    # Ground
    draw.rectangle([0, 500, 1280, 720], fill=(80, 140, 60))

    # Trees
    for tx in [100, 350, 700, 1050]:
        draw.rectangle([tx - 8, 420, tx + 8, 500], fill=(100, 70, 40))
        draw.polygon([(tx - 40, 500), (tx, 360), (tx + 40, 500)], fill=(30, 120, 30))
        draw.polygon([(tx - 35, 460), (tx, 330), (tx + 35, 460)], fill=(40, 140, 40))

    # River
    draw.polygon([(500, 720), (550, 500), (600, 500), (700, 720)], fill=(40, 100, 180))

    # Clouds
    for cx, cy in [(200, 80), (600, 60), (1100, 100)]:
        draw.ellipse([cx - 50, cy - 20, cx + 50, cy + 20], fill=(220, 230, 240))
        draw.ellipse([cx - 30, cy - 35, cx + 30, cy + 5], fill=(230, 240, 250))
        draw.ellipse([cx + 20, cy - 25, cx + 70, cy + 15], fill=(225, 235, 245))

    # Title text on image
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 36)
    except:
        font = ImageFont.load_default()
    draw.text((80, 580), "The Journey Begins - Chapter 3", fill=(255, 255, 240), font=font)

    img.save(BG_IMAGE)
    print(f"Illustration created: {BG_IMAGE}")


def create_initial():
    create_illustration()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # --- Slide 1: Title Slide ---
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "The Enchanted Forest"
    s1.placeholders[1].text = "An Interactive Story Presentation\nBy Elena Vasquez"
    fill = s1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    for para in s1.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.size = Pt(44)
    for para in s1.placeholders[1].text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    # --- Slide 2: Chapter 1 ---
    s2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    fill = s2.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    tb = s2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "Chapter 1: The Awakening"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xE8, 0xD4, 0x4D)
    body = s2.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(4))
    body.text_frame.word_wrap = True
    bp = body.text_frame.paragraphs[0]
    bp.text = ("Amara woke to the sound of wind chimes made from polished river stones. "
               "The morning light filtered through woven curtains, casting geometric shadows "
               "across the floor of her cottage. Today was different — she could feel it in the "
               "way the forest hummed, a low vibration that made the dust motes dance.")
    for r in bp.runs:
        r.font.size = Pt(18)
        r.font.color.rgb = RGBColor(0xEC, 0xF0, 0xF1)

    # --- Slide 3: Character Introduction ---
    s3 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = s3.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x34, 0x49, 0x5E)
    header = s3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    hp = header.text_frame.paragraphs[0]
    hp.text = "Characters"
    hp.alignment = PP_ALIGN.CENTER
    for r in hp.runs:
        r.font.size = Pt(32)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    chars = [
        ("Amara Osei", "The Keeper of Stories", "A young woman who discovered she can step into illustrations."),
        ("Kael Thornwood", "The Wandering Cartographer", "Maps the boundaries between the real and illustrated worlds."),
        ("Professor Linden", "The Archivist", "Protects ancient books whose illustrations hold living worlds."),
    ]
    for i, (name, title, desc) in enumerate(chars):
        y = Inches(1.8 + i * 1.8)
        tb = s3.shapes.add_textbox(Inches(1), y, Inches(11), Inches(1.5))
        tb.text_frame.word_wrap = True
        np = tb.text_frame.paragraphs[0]
        np.text = name
        for r in np.runs:
            r.font.size = Pt(22)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0xE8, 0xD4, 0x4D)
        tp = tb.text_frame.add_paragraph()
        tp.text = f"{title} — {desc}"
        for r in tp.runs:
            r.font.size = Pt(16)
            r.font.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)

    # --- Slide 4: Map ---
    s4 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = s4.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1C, 0x2E, 0x40)
    tb = s4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "The Illustrated Realms"
    p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.size = Pt(30)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # Add location markers as shapes
    locations = [
        ("Whispering Woods", Inches(2), Inches(2)),
        ("Crystal Caverns", Inches(7), Inches(3)),
        ("Sky Citadel", Inches(4), Inches(1.5)),
        ("Ember Plains", Inches(9), Inches(5)),
        ("Tidepool Gardens", Inches(1.5), Inches(5)),
    ]
    for loc_name, x, y in locations:
        shape = s4.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.4), Inches(0.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xE8, 0xD4, 0x4D)
        shape.line.fill.background()
        lbl = s4.shapes.add_textbox(x - Inches(0.5), y + Inches(0.45), Inches(2), Inches(0.5))
        lp = lbl.text_frame.paragraphs[0]
        lp.text = loc_name
        for r in lp.runs:
            r.font.size = Pt(12)
            r.font.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)

    # --- Slide 5: Chapter 2 ---
    s5 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = s5.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    tb = s5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "Chapter 2: Into the Pages"
    p.alignment = PP_ALIGN.LEFT
    for r in p.runs:
        r.font.size = Pt(36)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xE8, 0xD4, 0x4D)
    body = s5.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(4))
    body.text_frame.word_wrap = True
    bp = body.text_frame.paragraphs[0]
    bp.text = ("The book lay open on the pedestal, its pages shimmering with an inner light. "
               "Amara reached out, her fingers hovering above the illustration of the Whispering "
               "Woods. The trees in the drawing seemed to sway, and she could almost hear the "
               "rustle of leaves, the call of birds she had never seen.")
    for r in bp.runs:
        r.font.size = Pt(18)
        r.font.color.rgb = RGBColor(0xEC, 0xF0, 0xF1)

    # --- Slide 6: Quote ---
    s6 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = s6.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    tb = s6.shapes.add_textbox(Inches(2), Inches(2), Inches(9), Inches(3))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = '"Every story is a door. The question is whether you have the courage to step through."'
    p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.size = Pt(28)
        r.font.italic = True
        r.font.color.rgb = RGBColor(0xE8, 0xD4, 0x4D)
    attr = tb.text_frame.add_paragraph()
    attr.text = "— Professor Linden, The Archivist"
    attr.alignment = PP_ALIGN.RIGHT
    for r in attr.runs:
        r.font.size = Pt(18)
        r.font.color.rgb = RGBColor(0x95, 0xA5, 0xA6)

    # --- Slide 7: Chapter 3 Intro ---
    s7 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = s7.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    tb = s7.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "Chapter 3: The Page Turn"
    p.alignment = PP_ALIGN.LEFT
    for r in p.runs:
        r.font.size = Pt(36)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xE8, 0xD4, 0x4D)
    body = s7.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(4))
    body.text_frame.word_wrap = True
    bp = body.text_frame.paragraphs[0]
    bp.text = ("Kael spread the ancient map across the stone table. 'The next realm lies "
               "beyond this page,' he said, tracing a path through illustrated mountains. "
               "'When the page turns, the story reveals itself — but only to those who are "
               "watching.' Amara nodded, ready for what came next.")
    for r in bp.runs:
        r.font.size = Pt(18)
        r.font.color.rgb = RGBColor(0xEC, 0xF0, 0xF1)

    # --- Slide 8: The Reveal Slide (KEY SLIDE) ---
    s8 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    # 1. Background illustration image (full slide)
    pic = s8.shapes.add_picture(BG_IMAGE, 0, 0, slide_w, slide_h)

    # 2. White rectangle overlay covering entire slide (on top of illustration)
    white_rect = s8.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, slide_w, slide_h
    )
    white_rect.fill.solid()
    white_rect.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    white_rect.line.fill.background()  # no border
    # Name it for identification
    white_rect.name = "WhiteOverlay"

    # 3. Narration text box (positioned on slide, visually behind the white rect)
    narration = s8.shapes.add_textbox(Inches(1.5), Inches(5), Inches(10), Inches(2))
    narration.text_frame.word_wrap = True
    narration.name = "NarrationBox"
    np = narration.text_frame.paragraphs[0]
    np.text = ("As the white veil lifted, the enchanted forest materialized in vivid color. "
               "Amara stood at the edge of the illustrated world, the painted trees towering "
               "above her, their leaves whispering secrets of forgotten tales. The journey "
               "through Chapter 3 had truly begun.")
    np.alignment = PP_ALIGN.CENTER
    for r in np.runs:
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xE0)

    # --- Slide 9: Reflection ---
    s9 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = s9.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    tb = s9.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "Reflections"
    p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.size = Pt(36)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    body = s9.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(4))
    body.text_frame.word_wrap = True
    bp = body.text_frame.paragraphs[0]
    bp.text = ("The illustrated world challenged everything Amara thought she knew about "
               "reality. Each page turn revealed not just new landscapes, but new truths "
               "about the nature of stories themselves. Were they merely ink on paper, or "
               "something far more alive?")
    for r in bp.runs:
        r.font.size = Pt(18)
        r.font.color.rgb = RGBColor(0xEC, 0xF0, 0xF1)

    # --- Slide 10: To Be Continued ---
    s10 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = s10.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    tb = s10.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = "To Be Continued..."
    p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.size = Pt(48)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xE8, 0xD4, 0x4D)
    sub = tb.text_frame.add_paragraph()
    sub.text = "The next chapter awaits your click."
    sub.alignment = PP_ALIGN.CENTER
    for r in sub.runs:
        r.font.size = Pt(20)
        r.font.color.rgb = RGBColor(0x95, 0xA5, 0xA6)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
