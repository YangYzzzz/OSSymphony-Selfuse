"""
Initial Setup: Brand identity presentation - 8 slides, slide 2 has 3 textboxes with incorrect font sizes
Task ID: osworld_impress_textbox_fontsize_specific_006
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_textbox_fontsize_specific_006'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_layouts = prs.slide_layouts

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(slide_layouts[0])
    slide1.shapes.title.text = "Horizon Creative Studio"
    slide1.placeholders[1].text = "Brand Identity Guide 2025"
    fill1 = slide1.background.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    for ph in [slide1.shapes.title, slide1.placeholders[1]]:
        for para in ph.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # ---- Slide 2: Typography Scale (the key slide) ----
    # Three textboxes with sizes that are NOT the target (72, 36, 18)
    # Using 48pt, 24pt, 12pt instead
    slide2 = prs.slides.add_slide(slide_layouts[5])  # blank
    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF0)

    # First textbox: title-level text, size 48pt (target: 72pt)
    tb1 = slide2.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(10.0), Inches(1.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = False
    p1 = tf1.paragraphs[0]
    p1.text = "Horizon Creative Studio"
    r1 = p1.runs[0]
    r1.font.name = "Calibri"
    r1.font.size = Pt(48)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # Second textbox: subtitle-level text, size 24pt (target: 36pt)
    tb2 = slide2.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(10.0), Inches(1.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = False
    p2 = tf2.paragraphs[0]
    p2.text = "Brand Identity & Visual Language"
    r2 = p2.runs[0]
    r2.font.name = "Calibri"
    r2.font.size = Pt(24)
    r2.font.bold = False
    r2.font.color.rgb = RGBColor(0x4A, 0x4A, 0x6A)

    # Third textbox: tagline-level text, size 12pt (target: 18pt)
    tb3 = slide2.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(10.0), Inches(1.0))
    tf3 = tb3.text_frame
    tf3.word_wrap = False
    p3 = tf3.paragraphs[0]
    p3.text = "Crafting Authentic Connections Through Design"
    r3 = p3.runs[0]
    r3.font.name = "Calibri"
    r3.font.size = Pt(12)
    r3.font.italic = True
    r3.font.color.rgb = RGBColor(0x8A, 0x8A, 0xAA)

    # ---- Slide 3: Brand Colors ----
    slide3 = prs.slides.add_slide(slide_layouts[1])
    slide3.shapes.title.text = "Brand Color Palette"
    slide3.placeholders[1].text = (
        "Primary: Deep Navy #1A1A2E\n"
        "Secondary: Electric Teal #0F9B8E\n"
        "Accent: Warm Amber #F5A623\n"
        "Neutral: Soft Ivory #F5F5F0\n"
        "Text: Charcoal #3D3D3D"
    )

    # ---- Slide 4: Typography Guide ----
    slide4 = prs.slides.add_slide(slide_layouts[1])
    slide4.shapes.title.text = "Typography System"
    slide4.placeholders[1].text = (
        "Primary Typeface: Calibri\n"
        "Headlines: Bold, Navy\n"
        "Body Text: Regular, Charcoal\n"
        "Captions: Italic, Muted Navy"
    )

    # ---- Slide 5: Logo Usage ----
    slide5 = prs.slides.add_slide(slide_layouts[1])
    slide5.shapes.title.text = "Logo Usage Guidelines"
    slide5.placeholders[1].text = (
        "Clear Space: Minimum 20px on all sides\n"
        "Minimum Size: 80px wide for digital, 25mm for print\n"
        "Approved Backgrounds: White, Navy, Teal\n"
        "Never: Stretch, rotate, or recolor the logo"
    )

    # ---- Slide 6: Photography Style ----
    slide6 = prs.slides.add_slide(slide_layouts[1])
    slide6.shapes.title.text = "Photography & Imagery"
    slide6.placeholders[1].text = (
        "Style: Authentic, warm-toned lifestyle photography\n"
        "Subjects: Real people in real environments\n"
        "Color Grading: Warm highlights, deep shadows\n"
        "Avoid: Stock-looking, overly staged images"
    )

    # ---- Slide 7: Tone of Voice ----
    slide7 = prs.slides.add_slide(slide_layouts[1])
    slide7.shapes.title.text = "Tone of Voice"
    slide7.placeholders[1].text = (
        "Confident but not arrogant\n"
        "Creative but not chaotic\n"
        "Human but not casual\n"
        "Inspiring but not overwhelming"
    )

    # ---- Slide 8: Contact & Resources ----
    slide8 = prs.slides.add_slide(slide_layouts[0])
    slide8.shapes.title.text = "Questions & Resources"
    slide8.placeholders[1].text = "brand@horizoncreative.studio  |  horizoncreative.studio/brand-guide"
    fill8 = slide8.background.fill
    fill8.solid()
    fill8.fore_color.rgb = RGBColor(0x0F, 0x9B, 0x8E)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
