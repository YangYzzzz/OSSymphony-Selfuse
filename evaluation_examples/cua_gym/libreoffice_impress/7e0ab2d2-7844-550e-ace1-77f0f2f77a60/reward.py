"""
Reward Script: Set font sizes on slide 2 textboxes to 72pt, 36pt, 18pt
Task ID: osworld_impress_textbox_fontsize_specific_006
Domain: libreoffice_impress
Scoring:
  Component 1: First textbox (TextBox 2) font size == 72pt      — 0.4 pts
  Component 2: Second textbox (TextBox 3) font size == 36pt     — 0.3 pts
  Component 3: Third textbox (TextBox 4) font size == 18pt      — 0.3 pts
  Total: 1.0
"""

import os
from pptx import Presentation
from pptx.util import Pt

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_textbox_fontsize_specific_006'

# Target font sizes in EMU (1 pt = 12700 EMU)
# 72pt = 914400, 36pt = 457200, 18pt = 228600
TARGET_SIZES_PT = [72.0, 36.0, 18.0]
TARGET_SIZES_EMU = [int(pt * 12700) for pt in TARGET_SIZES_PT]


def get_textboxes_on_slide2(prs):
    """
    Return the three textbox shapes on slide 2 (index 1) in order.
    Skips 'Title 1' placeholder (shape index 0), returns TextBox 2, 3, 4.
    """
    slide2 = prs.slides[1]
    textboxes = []
    for shape in slide2.shapes:
        # Include only non-title text-frame shapes (TextBox 2, 3, 4)
        if shape.has_text_frame and shape.name != 'Title 1':
            textboxes.append(shape)
    return textboxes


def get_shape_font_size_pt(shape):
    """
    Return the font size in points for the first non-empty run in the shape.
    Checks all paragraphs and runs. Returns None if no explicit size found.
    """
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                return run.font.size / 12700.0
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: verify slide 2 exists and has 3 textboxes
    try:
        if len(prs.slides) < 2:
            print("CRITICAL: Presentation has fewer than 2 slides")
            print("REWARD: 0.0")
            return 0.0

        textboxes = get_textboxes_on_slide2(prs)
        if len(textboxes) < 3:
            print(f"CRITICAL: Expected 3 textboxes on slide 2, found {len(textboxes)}")
            print("REWARD: 0.0")
            return 0.0
    except Exception as e:
        print(f"CRITICAL: Cannot inspect slide 2: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: First textbox (TextBox 2) font size == 72pt (0.4 points)
    try:
        tb1 = textboxes[0]
        size_pt = get_shape_font_size_pt(tb1)
        expected_pt = TARGET_SIZES_PT[0]  # 72pt
        if size_pt is not None and abs(size_pt - expected_pt) < 0.1:
            print(f"PASS: Component 1 — TextBox 2 ('{tb1.name}') font size is {size_pt}pt == {expected_pt}pt (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 — TextBox 2 ('{tb1.name}') font size is {size_pt}pt, expected {expected_pt}pt")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Second textbox (TextBox 3) font size == 36pt (0.3 points)
    try:
        tb2 = textboxes[1]
        size_pt = get_shape_font_size_pt(tb2)
        expected_pt = TARGET_SIZES_PT[1]  # 36pt
        if size_pt is not None and abs(size_pt - expected_pt) < 0.1:
            print(f"PASS: Component 2 — TextBox 3 ('{tb2.name}') font size is {size_pt}pt == {expected_pt}pt (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 — TextBox 3 ('{tb2.name}') font size is {size_pt}pt, expected {expected_pt}pt")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Third textbox (TextBox 4) font size == 18pt (0.3 points)
    try:
        tb3 = textboxes[2]
        size_pt = get_shape_font_size_pt(tb3)
        expected_pt = TARGET_SIZES_PT[2]  # 18pt
        if size_pt is not None and abs(size_pt - expected_pt) < 0.1:
            print(f"PASS: Component 3 — TextBox 4 ('{tb3.name}') font size is {size_pt}pt == {expected_pt}pt (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 3 — TextBox 4 ('{tb3.name}') font size is {size_pt}pt, expected {expected_pt}pt")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point: run verification against the canonical artifact path on the VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
