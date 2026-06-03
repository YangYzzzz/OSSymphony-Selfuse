"""
Initial Setup: Create a 10-slide Roadmap presentation with slide 7 titled 'Strategic Roadmap' but blank.
Task ID: impress_exec_014
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_014'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.text = body_lines[0]
    for line in body_lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
    return slide


def add_blank_titled_slide(prs, title_text):
    """Add a slide with only a title (Title Only layout, index 5)."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title_text
    return slide


def create_initial():
    prs = Presentation()
    # Standard widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs, "NovaTech Inc. — Strategic Plan 2025-2026",
                    "Board Review Presentation | Confidential")

    # Slide 2: Company Overview
    add_content_slide(prs, "Company Overview", [
        "Founded in 2018 with a mission to modernize supply-chain analytics",
        "Headquarters in Austin, TX with offices in London and Singapore",
        "450+ employees across engineering, sales, and operations",
        "Annual recurring revenue: $38.2M (FY2024)",
        "Key verticals: retail, logistics, manufacturing",
    ])

    # Slide 3: Market Opportunity
    add_content_slide(prs, "Market Opportunity", [
        "Total addressable market: $12.4B by 2027 (Gartner)",
        "Current market penetration: ~0.3%",
        "Top competitor landscape: SAP Ariba, Coupa, Jaggaer",
        "Our differentiator: real-time predictive analytics engine",
        "Customer NPS score: 72 (industry avg: 48)",
    ])

    # Slide 4: Financial Performance
    add_content_slide(prs, "Financial Performance — FY2024", [
        "Revenue: $38.2M (+42% YoY)",
        "Gross margin: 78%",
        "Operating expenses: $29.1M",
        "EBITDA: $4.8M (first positive year)",
        "Cash runway: 24 months at current burn rate",
        "Series C raised: $65M in March 2024",
    ])

    # Slide 5: Product Roadmap
    add_content_slide(prs, "Product Roadmap", [
        "Q1 2025: Launch AI-powered demand forecasting module",
        "Q2 2025: Expand to 3 new geographic markets (DACH, ANZ, Brazil)",
        "Q3 2025: Strategic partnership with major logistics provider",
        "Q4 2025: Begin IPO preparation and SOX compliance audit",
        "Q1 2026: Target IPO on NASDAQ",
    ])

    # Slide 6: Team & Leadership
    add_content_slide(prs, "Leadership Team", [
        "CEO — Rachel Torres (ex-VP at Oracle Cloud)",
        "CTO — David Kim (ex-Principal Engineer at AWS)",
        "CFO — Amanda Osei (ex-Director at Goldman Sachs)",
        "VP Sales — Carlos Rivera (ex-Regional Director at Salesforce)",
        "VP Engineering — Priya Sharma (ex-Staff at Google)",
    ])

    # Slide 7: Strategic Roadmap — BLANK (this is the task slide)
    add_blank_titled_slide(prs, "Strategic Roadmap")

    # Slide 8: Risk Analysis
    add_content_slide(prs, "Risk Analysis", [
        "Market risk: economic downturn could slow enterprise adoption",
        "Competition: incumbents accelerating AI feature development",
        "Talent retention: 12% attrition rate in engineering",
        "Regulatory: data privacy laws in EU and Brazil evolving",
        "Mitigation: diversified verticals, competitive comp, legal counsel",
    ])

    # Slide 9: Key Metrics & KPIs
    add_content_slide(prs, "Key Metrics & KPIs", [
        "Monthly active users: 12,400 (+18% QoQ)",
        "Customer acquisition cost: $2,840",
        "Lifetime value: $42,600 (LTV/CAC = 15x)",
        "Churn rate: 3.2% annual",
        "Average contract value: $86,000",
    ])

    # Slide 10: Next Steps
    add_content_slide(prs, "Next Steps & Timeline", [
        "Board approval of IPO timeline by April 15, 2025",
        "Engage underwriting banks by May 2025",
        "Complete SOX compliance audit by September 2025",
        "File S-1 prospectus by December 2025",
        "Target IPO date: Q1 2026",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
