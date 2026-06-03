"""
Reward Script: Timeline on slide 7 with 5 milestones
Task ID: impress_exec_014
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Horizontal line at ~y=4in on slide 7
  Component 2 (0.30): 5 oval/circle shapes filled #003366
  Component 3 (0.25): 5 correct milestone text labels
  Component 4 (0.20): Circle positions at approximately correct x-coordinates
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_014'

# Expected milestones with approximate x positions (inches)
MILESTONES = [
    ('Q1 Launch', 1.0),
    ('Q2 Expansion', 3.5),
    ('Q3 Partnership', 6.0),
    ('Q4 IPO Prep', 8.5),
    ('Q1 2026 IPO', 11.0),
]

def emu_to_inches(emu):
    return emu / 914400.0

def approx_equal(val1, val2, tolerance_in=0.6):
    """Check if two inch values are approximately equal within tolerance."""
    return abs(val1 - val2) <= tolerance_in

def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 7 slides
    if len(prs.slides) < 7:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 7")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[6]  # 0-indexed, slide 7

    # Collect non-placeholder shapes on slide 7 (these are the task-introduced shapes)
    new_shapes = [s for s in slide.shapes if s.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER]

    # Component 1: Horizontal line at ~y=4in (0.25 points)
    # The line is implemented as a thin rectangle spanning across the slide
    try:
        line_found = False
        for shape in new_shapes:
            top_in = emu_to_inches(shape.top)
            h_in = emu_to_inches(shape.height)
            w_in = emu_to_inches(shape.width)
            # A line-like shape: very thin height (<0.1in), wide (>5in), near y=4in
            if h_in < 0.15 and w_in > 5.0 and approx_equal(top_in, 4.0, 0.5):
                line_found = True
                print(f"PASS: Component 1 — Horizontal line found at y={top_in:.2f}in, width={w_in:.2f}in, height={h_in:.3f}in (0.25 pts)")
                break
        if line_found:
            total_score += 0.25
        else:
            # Also check for freeform/connector line shapes
            for shape in new_shapes:
                if shape.shape_type in (MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM):
                    top_in = emu_to_inches(shape.top)
                    w_in = emu_to_inches(shape.width)
                    if w_in > 5.0 and approx_equal(top_in, 4.0, 0.5):
                        line_found = True
                        print(f"PASS: Component 1 — Line shape found at y={top_in:.2f}in (0.25 pts)")
                        total_score += 0.25
                        break
            if not line_found:
                print(f"FAIL: Component 1 — No horizontal line found near y=4in on slide 7")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: 5 oval/circle shapes filled #003366 (0.30 points)
    try:
        circles = []
        for shape in new_shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                try:
                    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
                    if shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.OVAL:
                        circles.append(shape)
                        continue
                except:
                    pass
            # Also accept ovals by checking near-square dimensions near 0.4in
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                w_in = emu_to_inches(shape.width)
                h_in = emu_to_inches(shape.height)
                if abs(w_in - h_in) < 0.1 and 0.2 <= w_in <= 0.8:
                    circles.append(shape)

        # Deduplicate (some may be added twice)
        seen_ids = set()
        unique_circles = []
        for c in circles:
            sid = id(c)
            if sid not in seen_ids:
                seen_ids.add(sid)
                unique_circles.append(c)
        circles = unique_circles

        # Check fill color
        filled_circles = []
        for c in circles:
            try:
                fill = c.fill
                if fill.type is not None and fill.type == 1:  # SOLID
                    color = str(fill.fore_color.rgb).upper()
                    if color == '003366':
                        filled_circles.append(c)
                    else:
                        print(f"  INFO: Circle at x={emu_to_inches(c.left):.2f}in has fill color {color}, expected 003366")
                        filled_circles.append(c)  # Still count it but note the color
            except:
                filled_circles.append(c)  # Count it even if color check fails

        num_correct_color = sum(
            1 for c in circles
            if _check_fill_003366(c)
        )

        if len(circles) >= 5 and num_correct_color >= 5:
            print(f"PASS: Component 2 — Found {len(circles)} circles, {num_correct_color} with #003366 fill (0.30 pts)")
            total_score += 0.30
        elif len(circles) >= 5:
            # Partial: circles exist but color may be wrong
            pts = 0.15
            print(f"PARTIAL: Component 2 — Found {len(circles)} circles but only {num_correct_color} with #003366 fill ({pts} pts)")
            total_score += pts
        elif len(circles) >= 3:
            pts = 0.10
            print(f"PARTIAL: Component 2 — Found only {len(circles)} circles (expected 5) ({pts} pts)")
            total_score += pts
        else:
            print(f"FAIL: Component 2 — Found only {len(circles)} circles (expected 5)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: 5 correct milestone text labels (0.25 points)
    try:
        expected_labels = {'Q1 Launch', 'Q2 Expansion', 'Q3 Partnership', 'Q4 IPO Prep', 'Q1 2026 IPO'}
        found_labels = set()
        for shape in new_shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text in expected_labels:
                    found_labels.add(text)

        matched = len(found_labels)
        if matched == 5:
            print(f"PASS: Component 3 — All 5 milestone labels found: {found_labels} (0.25 pts)")
            total_score += 0.25
        elif matched >= 3:
            pts = round(0.25 * matched / 5, 2)
            print(f"PARTIAL: Component 3 — Found {matched}/5 labels: {found_labels} ({pts} pts)")
            total_score += pts
        else:
            print(f"FAIL: Component 3 — Found only {matched}/5 labels: {found_labels}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Circles at approximately correct x positions (0.20 points)
    try:
        # Map circle centers to expected positions
        circle_centers = []
        for c in circles:
            center_x = emu_to_inches(c.left) + emu_to_inches(c.width) / 2
            circle_centers.append(center_x)

        expected_xs = [m[1] for m in MILESTONES]  # 1.0, 3.5, 6.0, 8.5, 11.0

        # For each expected x, find the closest circle center
        matched_positions = 0
        for ex in expected_xs:
            best_dist = float('inf')
            for cx in circle_centers:
                dist = abs(cx - ex)
                if dist < best_dist:
                    best_dist = dist
            if best_dist <= 0.6:  # within 0.6 inch tolerance
                matched_positions += 1

        if matched_positions == 5:
            print(f"PASS: Component 4 — All 5 circles at correct x positions (0.20 pts)")
            total_score += 0.20
        elif matched_positions >= 3:
            pts = round(0.20 * matched_positions / 5, 2)
            print(f"PARTIAL: Component 4 — {matched_positions}/5 circles at correct positions ({pts} pts)")
            total_score += pts
        else:
            print(f"FAIL: Component 4 — Only {matched_positions}/5 circles at correct positions")
            print(f"  Circle centers: {[f'{x:.2f}' for x in sorted(circle_centers)]}")
            print(f"  Expected: {expected_xs}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


def _check_fill_003366(shape):
    """Check if shape has solid fill with color #003366."""
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:
            return str(fill.fore_color.rgb).upper() == '003366'
    except:
        pass
    return False


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
