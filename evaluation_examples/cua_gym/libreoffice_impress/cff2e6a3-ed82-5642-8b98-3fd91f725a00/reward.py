"""
Reward Script: Resize the image on slide 6 to be as large as possible while
               fitting within the slide bounds and keeping it proportional.
Task ID: osworld_impress_image_fill_slide_012
Domain: libreoffice_impress
Scoring:
  Component 1 — Image size reaches maximum proportional fit (0.5 pts)
  Component 2 — Image aspect ratio is preserved after resize (0.3 pts)
  Component 3 — Image fits within slide bounds (no overflow) (0.2 pts)
"""

import os
import io

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_image_fill_slide_012'

def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0

    Task: Resize the image on slide 6 to fill the slide bounds proportionally.

    The image has pixel dims 640x480 (4:3 aspect ratio).
    The slide is 10 x 7.5 inches = 9144000 x 6858000 EMU (also 4:3).
    Maximum proportional fit = full slide: 9144000 x 6858000 EMU.
    """
    total_score = 0.0

    # --- Load the presentation ---
    try:
        from pptx import Presentation
        from PIL import Image
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # --- Precondition: Slide 6 must exist and have a picture ---
    try:
        if len(prs.slides) < 6:
            print(f"FAIL: Presentation has fewer than 6 slides ({len(prs.slides)} slides)")
            print("REWARD: 0.0")
            return 0.0

        slide = prs.slides[5]  # slide 6 (0-indexed)
        slide_width = prs.slide_width    # 9144000 EMU = 10 inches
        slide_height = prs.slide_height  # 6858000 EMU = 7.5 inches

        # Find the picture shape on slide 6
        picture_shape = None
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                picture_shape = shape
                break

        if picture_shape is None:
            print("FAIL: No picture shape found on slide 6")
            print("REWARD: 0.0")
            return 0.0

    except Exception as e:
        print(f"CRITICAL: Error accessing slide 6 or picture: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Gather actual dimensions
    actual_width = picture_shape.width
    actual_height = picture_shape.height
    actual_left = picture_shape.left
    actual_top = picture_shape.top

    print(f"INFO: Slide dimensions: {slide_width} x {slide_height} EMU ({slide_width/914400:.4f} x {slide_height/914400:.4f} inches)")
    print(f"INFO: Picture dimensions: {actual_width} x {actual_height} EMU ({actual_width/914400:.4f} x {actual_height/914400:.4f} inches)")
    print(f"INFO: Picture position: left={actual_left}, top={actual_top}")

    # Determine the expected maximum proportional size based on image pixel dimensions
    try:
        img_blob = picture_shape.image.blob
        img = Image.open(io.BytesIO(img_blob))
        pixel_w, pixel_h = img.size
        image_aspect_ratio = pixel_w / pixel_h
        slide_aspect_ratio = slide_width / slide_height

        # Compute max-fit proportional dimensions
        if image_aspect_ratio >= slide_aspect_ratio:
            # Image is wider relative to slide: width-constrained
            expected_max_width = slide_width
            expected_max_height = int(round(slide_width / image_aspect_ratio))
        else:
            # Image is taller relative to slide: height-constrained
            expected_max_height = slide_height
            expected_max_width = int(round(slide_height * image_aspect_ratio))

        print(f"INFO: Image pixel size: {pixel_w} x {pixel_h}, aspect ratio: {image_aspect_ratio:.6f}")
        print(f"INFO: Expected max proportional size: {expected_max_width} x {expected_max_height} EMU")
    except Exception as e:
        # If PIL is unavailable, fallback to slide dimensions as expected size
        # (for 4:3 image on 4:3 slide, the answer is the full slide)
        expected_max_width = slide_width
        expected_max_height = slide_height
        print(f"WARN: Could not determine image pixel dims via PIL: {e}. Using slide size as expected max.")

    # Tolerance: 0.5% relative tolerance (consistent with SKILL.md guidelines)
    def within_tolerance(val, expected, tol=0.005):
        if expected == 0:
            return val == 0
        return abs(val - expected) / max(abs(val), abs(expected)) <= tol

    # --- Component 1: Image size reaches maximum proportional fit (0.5 points) ---
    # This FAILS on initial (image is 3" x 2.25") and PASSES on golden (image is 10" x 7.5").
    try:
        width_ok = within_tolerance(actual_width, expected_max_width)
        height_ok = within_tolerance(actual_height, expected_max_height)

        if width_ok and height_ok:
            print(f"PASS: Component 1 — Image resized to maximum proportional fit: "
                  f"{actual_width} x {actual_height} EMU "
                  f"(expected ~{expected_max_width} x {expected_max_height} EMU) (0.5 pts)")
            total_score += 0.5
        else:
            # Partial credit: if one dimension reaches ~80% of maximum proportional extent
            width_fraction = actual_width / expected_max_width if expected_max_width > 0 else 0
            height_fraction = actual_height / expected_max_height if expected_max_height > 0 else 0
            avg_fraction = (width_fraction + height_fraction) / 2
            if avg_fraction >= 0.8:
                print(f"PARTIAL: Component 1 — Image resized but not to maximum: "
                      f"actual={actual_width} x {actual_height}, "
                      f"expected ~{expected_max_width} x {expected_max_height} "
                      f"(coverage: {avg_fraction:.1%}) (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 1 — Image not resized to maximum: "
                      f"actual={actual_width} x {actual_height} EMU "
                      f"({actual_width/914400:.4f} x {actual_height/914400:.4f} in), "
                      f"expected ~{expected_max_width} x {expected_max_height} EMU")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # --- Component 2: Image aspect ratio preserved after resize (0.3 points) ---
    # The image aspect ratio should match the original pixel aspect ratio.
    # On initial: image is already 4:3 (3.0 x 2.25), so this would pass — BUT we make it
    # a sub-condition of Component 1 criteria by also checking that the size changed significantly.
    # We score this only when the image was actually resized (width > initial width).
    #
    # Initial width: 2743200 EMU (3 inches). Golden width: 9144000 EMU (10 inches).
    # We only score proportionality if the image was actually enlarged beyond initial size.
    INITIAL_WIDTH_EMU = 2743200  # 3 inches — the known initial image width

    try:
        if actual_width <= INITIAL_WIDTH_EMU:
            print(f"FAIL: Component 2 — Image not enlarged beyond initial size "
                  f"(actual width {actual_width} <= initial {INITIAL_WIDTH_EMU}). Skipping ratio check.")
        else:
            actual_aspect = actual_width / actual_height if actual_height != 0 else 0

            # The expected aspect ratio is from the image pixel data (4:3 = 1.333...)
            try:
                img_blob = picture_shape.image.blob
                img2 = Image.open(io.BytesIO(img_blob))
                px_w, px_h = img2.size
                original_aspect = px_w / px_h
            except Exception:
                original_aspect = 4.0 / 3.0  # fallback: known 640x480 = 4:3

            aspect_tolerance = 0.01  # 1% tolerance on aspect ratio
            ratio_ok = abs(actual_aspect - original_aspect) / original_aspect <= aspect_tolerance

            if ratio_ok:
                print(f"PASS: Component 2 — Aspect ratio preserved after resize: "
                      f"actual={actual_aspect:.6f}, expected={original_aspect:.6f} (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 2 — Aspect ratio distorted: "
                      f"actual={actual_aspect:.6f}, expected={original_aspect:.6f} "
                      f"(deviation: {abs(actual_aspect - original_aspect)/original_aspect:.2%})")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # --- Component 3: Image fits within slide bounds (0.2 points) ---
    # Check that the image (left + width, top + height) does not exceed slide dimensions.
    # On initial, the image is 3"x2.25" at (1.0", 1.2"), so right=4" < 10", bottom=3.45" < 7.5" — passes!
    # BUT we gate this on the image being enlarged (same gate as Component 2).
    try:
        if actual_width <= INITIAL_WIDTH_EMU:
            print(f"FAIL: Component 3 — Image not enlarged; bounds check only meaningful post-resize.")
        else:
            right_edge = (actual_left or 0) + actual_width
            bottom_edge = (actual_top or 0) + actual_height

            # Allow 1% tolerance beyond slide edge (rounding artifacts)
            right_ok = right_edge <= slide_width * 1.01
            bottom_ok = bottom_edge <= slide_height * 1.01

            if right_ok and bottom_ok:
                print(f"PASS: Component 3 — Image fits within slide bounds: "
                      f"right_edge={right_edge/914400:.4f}\" <= {slide_width/914400:.4f}\", "
                      f"bottom_edge={bottom_edge/914400:.4f}\" <= {slide_height/914400:.4f}\" (0.2 pts)")
                total_score += 0.2
            else:
                overflow_details = []
                if not right_ok:
                    overflow_details.append(f"right overflow: {right_edge/914400:.4f}\" > {slide_width/914400:.4f}\"")
                if not bottom_ok:
                    overflow_details.append(f"bottom overflow: {bottom_edge/914400:.4f}\" > {slide_height/914400:.4f}\"")
                print(f"FAIL: Component 3 — Image overflows slide: {', '.join(overflow_details)}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
