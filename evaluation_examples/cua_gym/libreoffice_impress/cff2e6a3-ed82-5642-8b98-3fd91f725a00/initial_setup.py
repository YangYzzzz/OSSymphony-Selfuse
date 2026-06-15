"""
Initial Setup: Nature photography presentation with small wildlife photo on slide 6
Task ID: osworld_impress_image_fill_slide_012
Domain: libreoffice_impress
"""

import os
import io
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_image_fill_slide_012'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

# Slide dimensions (standard widescreen 10" x 7.5")
SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def make_nature_image(width_px, height_px, color_top, color_bottom, label=None):
    """Create a synthetic nature-like gradient image as bytes."""
    img = Image.new('RGB', (width_px, height_px))
    draw = ImageDraw.Draw(img)
    # Simple gradient
    for y in range(height_px):
        ratio = y / height_px
        r = int(color_top[0] * (1 - ratio) + color_bottom[0] * ratio)
        g = int(color_top[1] * (1 - ratio) + color_bottom[1] * ratio)
        b = int(color_top[2] * (1 - ratio) + color_bottom[2] * ratio)
        draw.line([(0, y), (width_px, y)], fill=(r, g, b))
    # Add some texture (random-looking dots)
    import random
    random.seed(42)
    for _ in range(300):
        x = random.randint(0, width_px - 1)
        y = random.randint(0, height_px - 1)
        draw.ellipse([x-2, y-2, x+2, y+2], fill=(
            min(255, r + random.randint(-30, 30)),
            min(255, g + random.randint(-30, 30)),
            min(255, b + random.randint(-30, 30))
        ))
    if label:
        draw.text((10, 10), label, fill=(255, 255, 255))
    buf = io.BytesIO()
    img.save(buf, format='PNG')
    buf.seek(0)
    return buf


def add_slide_bg(slide, prs, r, g, b):
    """Set slide background to a solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def add_title_text(slide, title, subtitle=None):
    """Add a title textbox to the slide."""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = 'Calibri'
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.name = 'Calibri'
        run2.font.size = Pt(20)
        run2.font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)


def create_initial():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # blank layout

    # ------- SLIDE 1: Title slide — "Wonders of the Wild" -------
    slide1 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide1, prs, 0x0D, 0x2B, 0x0D)  # dark forest green
    # Forest canopy image stretched across slide
    img1 = make_nature_image(800, 600, (34, 85, 34), (8, 40, 8))
    slide1.shapes.add_picture(img1, 0, 0, SLIDE_W, SLIDE_H)

    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Wonders of the Wild"
    run.font.name = 'Calibri'
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "A Journey Through Nature Photography"
    run2.font.name = 'Calibri'
    run2.font.size = Pt(22)
    run2.font.color.rgb = RGBColor(0xCC, 0xFF, 0xCC)

    # ------- SLIDE 2: "Majestic Landscapes" -------
    slide2 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide2, prs, 0x1A, 0x2E, 0x4A)  # deep blue sky
    img2 = make_nature_image(800, 400, (26, 92, 160), (139, 195, 74))
    slide2.shapes.add_picture(img2, Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    add_title_text(slide2, "Majestic Landscapes", "Mountains, Valleys & Plains")

    # ------- SLIDE 3: "Rainforest Life" -------
    slide3 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide3, prs, 0x00, 0x2B, 0x00)
    img3 = make_nature_image(600, 500, (0, 100, 0), (34, 139, 34))
    slide3.shapes.add_picture(img3, Inches(2), Inches(1.5), Inches(6), Inches(5))

    txt3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1.0))
    p3 = txt3.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = "Rainforest Life"
    r3.font.name = 'Calibri'
    r3.font.size = Pt(36)
    r3.font.bold = True
    r3.font.color.rgb = RGBColor(0xAD, 0xFF, 0x2F)

    cap3 = slide3.shapes.add_textbox(Inches(1), Inches(6.7), Inches(8), Inches(0.6))
    pc3 = cap3.text_frame.paragraphs[0]
    pc3.alignment = PP_ALIGN.CENTER
    rc3 = pc3.add_run()
    rc3.text = "The Amazon basin holds 10% of all species on Earth"
    rc3.font.name = 'Calibri'
    rc3.font.size = Pt(14)
    rc3.font.italic = True
    rc3.font.color.rgb = RGBColor(0xCC, 0xFF, 0x99)

    # ------- SLIDE 4: "Ocean Depths" -------
    slide4 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide4, prs, 0x00, 0x10, 0x40)
    img4 = make_nature_image(900, 500, (0, 20, 80), (0, 80, 160))
    slide4.shapes.add_picture(img4, Inches(0), Inches(1.2), Inches(10), Inches(5.5))
    add_title_text(slide4, "Ocean Depths", "Mysteries Beneath the Surface")

    cap4 = slide4.shapes.add_textbox(Inches(1), Inches(6.7), Inches(8), Inches(0.6))
    pc4 = cap4.text_frame.paragraphs[0]
    pc4.alignment = PP_ALIGN.LEFT
    rc4 = pc4.add_run()
    rc4.text = "Over 80% of the ocean remains unexplored by humans"
    rc4.font.name = 'Calibri'
    rc4.font.size = Pt(13)
    rc4.font.color.rgb = RGBColor(0x99, 0xDD, 0xFF)

    # ------- SLIDE 5: "Desert Solitude" -------
    slide5 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide5, prs, 0x3D, 0x2B, 0x00)
    img5 = make_nature_image(800, 480, (210, 140, 30), (150, 80, 10))
    slide5.shapes.add_picture(img5, Inches(0.3), Inches(1.3), Inches(9.4), Inches(5.5))
    add_title_text(slide5, "Desert Solitude", "Life in Arid Extremes")

    # ------- SLIDE 6: "Wildlife Encounters" — SMALL image in upper area -------
    slide6 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide6, prs, 0x1C, 0x1C, 0x1C)  # near black

    # Title
    txt6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.9))
    p6 = txt6.text_frame.paragraphs[0]
    p6.alignment = PP_ALIGN.CENTER
    r6 = p6.add_run()
    r6.text = "Wildlife Encounters"
    r6.font.name = 'Calibri'
    r6.font.size = Pt(32)
    r6.font.bold = True
    r6.font.color.rgb = RGBColor(0xFF, 0xCC, 0x00)

    # SMALL wildlife photo in the upper area (NOT maximized — this is what the agent must fix)
    # Image intrinsic size: 640x480 (4:3 ratio)
    # Placed small: 3" wide x 2.25" tall at position (1", 1.2") — clearly undersized
    img6 = make_nature_image(640, 480, (139, 90, 43), (80, 50, 10))
    slide6.shapes.add_picture(img6, Inches(1), Inches(1.2), Inches(3), Inches(2.25))

    cap6 = slide6.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.6))
    pc6 = cap6.text_frame.paragraphs[0]
    pc6.alignment = PP_ALIGN.CENTER
    rc6 = pc6.add_run()
    rc6.text = "Captured in the Serengeti, Tanzania"
    rc6.font.name = 'Calibri'
    rc6.font.size = Pt(15)
    rc6.font.italic = True
    rc6.font.color.rgb = RGBColor(0xCC, 0xAA, 0x55)

    note6 = slide6.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(2.5))
    tf6n = note6.text_frame
    tf6n.word_wrap = True
    pn = tf6n.paragraphs[0]
    pn.alignment = PP_ALIGN.LEFT
    rn = pn.add_run()
    rn.text = ("The Serengeti hosts the world's largest terrestrial mammal migration. "
               "Each year over 1.5 million wildebeest, 500,000 gazelles, and 250,000 "
               "zebras traverse the ecosystem following seasonal rains.")
    rn.font.name = 'Calibri'
    rn.font.size = Pt(13)
    rn.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

    # ------- SLIDE 7: "Conservation Matters" -------
    slide7 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide7, prs, 0x0A, 0x20, 0x0A)
    img7 = make_nature_image(700, 450, (20, 80, 20), (10, 50, 10))
    slide7.shapes.add_picture(img7, Inches(0.5), Inches(1.5), Inches(4), Inches(3.5))

    txt7 = slide7.shapes.add_textbox(Inches(5.0), Inches(1.2), Inches(4.5), Inches(5.5))
    tf7 = txt7.text_frame
    tf7.word_wrap = True
    ph7 = tf7.paragraphs[0]
    ph7.alignment = PP_ALIGN.LEFT
    rh7 = ph7.add_run()
    rh7.text = "Conservation Matters"
    rh7.font.name = 'Calibri'
    rh7.font.size = Pt(28)
    rh7.font.bold = True
    rh7.font.color.rgb = RGBColor(0x66, 0xFF, 0x66)

    facts = [
        "1 million species face extinction",
        "75% of land surface significantly altered",
        "66% of ocean area impacted by human activity",
        "We lose 4.7 million hectares of forest annually",
    ]
    for fact in facts:
        pb = tf7.add_paragraph()
        pb.alignment = PP_ALIGN.LEFT
        pb.level = 1
        rb = pb.add_run()
        rb.text = fact
        rb.font.name = 'Calibri'
        rb.font.size = Pt(14)
        rb.font.color.rgb = RGBColor(0xCC, 0xFF, 0xCC)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
