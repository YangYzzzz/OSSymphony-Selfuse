"""
Initial Setup: Create curriculum_overview.pptx on the VM Desktop
Task ID: impress_anim_079
Domain: libreoffice_impress

Creates a 5-slide curriculum overview presentation with no animations.
Slide 4 has a text box with 3 bullet points (no animation applied).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

WORKDIR = '/home/user/Desktop'
TASK_ID = 'impress_anim_079'
TASK_FILE = f'{WORKDIR}/curriculum_overview.pptx'
OUTPUT = f'/home/user/{TASK_ID}_initial.pptx'

def create_initial():
    prs = Presentation()

    # Slide dimensions: standard 10" x 7.5"
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ----------------------------------------------------------------
    # Slide 1: Title Slide
    # ----------------------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Secondary School Curriculum Overview"
    slide1.placeholders[1].text = "Academic Year 2025–2026\nDepartment of Curriculum & Instruction"

    # ----------------------------------------------------------------
    # Slide 2: Program Goals
    # ----------------------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Program Goals"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Develop critical thinking and analytical skills in all students"
    goals = [
        "Foster collaborative learning through project-based activities",
        "Integrate digital literacy across all subject areas",
        "Ensure equitable access to rigorous academic content",
        "Promote social-emotional learning and student wellbeing",
    ]
    for g in goals:
        p = tf2.add_paragraph()
        p.text = g
        p.level = 1

    # ----------------------------------------------------------------
    # Slide 3: Core Subject Areas
    # ----------------------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Core Subject Areas"
    tf3 = slide3.placeholders[1].text_frame
    subjects = [
        ("English Language Arts", "Reading, Writing, Speaking & Listening"),
        ("Mathematics", "Algebra, Geometry, Statistics & Data Analysis"),
        ("Science", "Biology, Chemistry, Physics, Environmental Science"),
        ("Social Studies", "World History, Civics, Economics, Geography"),
        ("World Languages", "Spanish, French, Mandarin Chinese"),
    ]
    tf3.text = subjects[0][0]
    tf3.paragraphs[0].runs[0].font.bold = True
    for subj, desc in subjects[1:]:
        p1 = tf3.add_paragraph()
        p1.text = subj
        p1.runs[0].font.bold = True
        p1.level = 0
        p2 = tf3.add_paragraph()
        p2.text = desc
        p2.level = 1

    # ----------------------------------------------------------------
    # Slide 4: Assessment Framework  (THE KEY SLIDE — no animations)
    # ----------------------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Assessment Framework"

    # Use the content placeholder for the 3 bullet points
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Formative assessments embedded in daily instruction to monitor student progress"
    bullets = [
        "Summative evaluations aligned to state standards at end of each unit",
        "Portfolio-based assessment capturing student growth over the semester",
    ]
    for b in bullets:
        p = tf4.add_paragraph()
        p.text = b
        p.level = 0

    # NO animations are added to this slide — that is the task for the agent

    # ----------------------------------------------------------------
    # Slide 5: Professional Development
    # ----------------------------------------------------------------
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Professional Development Plan"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Monthly collaborative planning sessions across grade levels and departments"
    pd_items = [
        "Annual summer institute focused on instructional best practices",
        "Peer observation and coaching cycles throughout the academic year",
        "Digital tools workshops: Google Classroom, Nearpod, Desmos",
    ]
    for item in pd_items:
        p = tf5.add_paragraph()
        p.text = item
        p.level = 1

    # Save - create the actual task file on the Desktop
    os.makedirs(WORKDIR, exist_ok=True)
    prs.save(TASK_FILE)

    # Also copy as the canonical initial tracking file
    import shutil
    shutil.copy(TASK_FILE, OUTPUT)

    print(f'Task file created: {TASK_FILE}')
    print(f'Initial tracking file created: {OUTPUT}')
    print('Slides:')
    print('  Slide 1: Title Slide')
    print('  Slide 2: Program Goals')
    print('  Slide 3: Core Subject Areas')
    print('  Slide 4: Assessment Framework (3 bullets, NO animation)')
    print('  Slide 5: Professional Development Plan')

create_initial()
