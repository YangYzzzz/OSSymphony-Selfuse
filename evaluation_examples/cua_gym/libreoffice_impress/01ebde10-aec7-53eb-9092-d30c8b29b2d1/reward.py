"""
Reward Script: Make last column ('Total') in slide 3 table bold with yellow (#FFF9C4) background
Task ID: impress_tct_028
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): All 7 cells in last column have yellow (#FFF9C4) background fill
  Component 2 (0.5): All 7 cells in last column have bold text (data rows changed from non-bold)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_028'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    # Find the table on slide 3 (index 2)
    slide = prs.slides[2]
    table = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            break

    if table is None:
        print("FAIL: No table found on slide 3")
        print("REWARD: 0.0")
        return 0.0

    num_rows = len(table.rows)
    num_cols = len(table.columns)

    # Precondition: table has expected dimensions
    if num_cols < 5 or num_rows < 7:
        print(f"FAIL: Table dimensions {num_rows}x{num_cols}, expected at least 7x5")
        print("REWARD: 0.0")
        return 0.0

    last_col = num_cols - 1

    # Precondition: last column header is 'Total'
    header_text = table.cell(0, last_col).text.strip()
    if header_text != 'Total':
        print(f"FAIL: Last column header is '{header_text}', expected 'Total'")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Yellow background fill (#FFF9C4) on all cells in last column (0.5 points)
    # In initial_env, all cells have FFFFFF. In golden_env, all should have FFF9C4.
    try:
        yellow_count = 0
        for r in range(num_rows):
            cell = table.cell(r, last_col)
            fill = cell.fill
            try:
                fg_rgb = str(fill.fore_color.rgb).upper()
            except Exception:
                fg_rgb = None

            if fg_rgb == 'FFF9C4':
                yellow_count += 1
            else:
                print(f"FAIL: Cell({r},{last_col}) background is {fg_rgb}, expected FFF9C4")

        if yellow_count == num_rows:
            print(f"PASS: Component 1 -- All {num_rows} cells in last column have #FFF9C4 background (0.5 pts)")
            total_score += 0.5
        elif yellow_count > 0:
            partial = 0.5 * (yellow_count / num_rows)
            print(f"PARTIAL: Component 1 -- {yellow_count}/{num_rows} cells have yellow background ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- No cells in last column have #FFF9C4 background")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Bold text on all cells in last column (0.5 points)
    # In initial_env, only header (row 0) is bold; data rows (1-6) are NOT bold.
    # In golden_env, ALL rows should be bold.
    # We only score the DATA rows (1-6) becoming bold, since header was already bold.
    try:
        bold_count = 0
        data_rows = num_rows - 1  # exclude header row which was already bold

        for r in range(1, num_rows):
            cell = table.cell(r, last_col)
            all_bold = True
            has_runs = False
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    has_runs = True
                    # None means inherit (treat as not bold), True means bold
                    if run.font.bold is not True:
                        all_bold = False

            if has_runs and all_bold:
                bold_count += 1
            else:
                print(f"FAIL: Cell({r},{last_col}) text is not bold")

        if bold_count == data_rows:
            print(f"PASS: Component 2 -- All {data_rows} data cells in last column are bold (0.5 pts)")
            total_score += 0.5
        elif bold_count > 0:
            partial = 0.5 * (bold_count / data_rows)
            print(f"PARTIAL: Component 2 -- {bold_count}/{data_rows} data cells are bold ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 -- No data cells in last column are bold")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
