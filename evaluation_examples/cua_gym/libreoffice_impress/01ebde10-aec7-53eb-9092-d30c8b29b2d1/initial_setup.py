"""
Initial Setup: Create Expense_Report.pptx with 5 slides, slide 3 has a table.
Task ID: impress_tct_028
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_028'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_cell_text(cell, text, font_name="Calibri", font_size=Pt(12),
                  bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to set cell text with formatting."""
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def set_cell_fill(cell, rgb_color):
    """Set cell background fill color."""
    tcPr = cell._tc.get_or_add_tcPr()
    solidFill = tcPr.makeelement(qn('a:solidFill'), {})
    srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': rgb_color})
    solidFill.append(srgbClr)
    # Remove any existing fill first
    for existing in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(existing)
    tcPr.append(solidFill)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Quarterly Expense Report"
    slide1.placeholders[1].text = "Fiscal Year 2025 — Prepared by Finance Department"

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    body2 = slide2.placeholders[1]
    tf2 = body2.text_frame
    tf2.text = "This report summarizes departmental expenses for Q1-Q3 of fiscal year 2025."
    p2 = tf2.add_paragraph()
    p2.text = "Key highlights include a 12% reduction in operational costs and a strategic reallocation of marketing budget toward digital channels."
    p3 = tf2.add_paragraph()
    p3.text = "Overall spending remains within the approved annual budget of $2.4M, with projected year-end savings of approximately $180K."

    # --- Slide 3: Table (5 columns x 7 rows) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Add title text box
    title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.6))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    run_title = p_title.add_run()
    run_title.text = "Departmental Expense Breakdown"
    run_title.font.name = "Calibri"
    run_title.font.size = Pt(24)
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    rows, cols = 7, 5
    table_shape = slide3.shapes.add_table(
        rows, cols,
        Inches(0.8), Inches(1.2),
        Inches(11.5), Inches(5.0)
    )
    table = table_shape.table

    # Set column widths
    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(2.0)
    table.columns[3].width = Inches(2.0)
    table.columns[4].width = Inches(2.5)

    # Headers
    headers = ["Department", "Q1", "Q2", "Q3", "Total"]
    header_color = RGBColor(0x33, 0x33, 0x33)
    for c, h in enumerate(headers):
        set_cell_text(table.cell(0, c), h, font_size=Pt(14), bold=True,
                      color=header_color, alignment=PP_ALIGN.CENTER)
        set_cell_fill(table.cell(0, c), "FFFFFF")

    # Data rows — all white background, non-bold
    data = [
        ["Engineering",   "$145,200", "$138,750", "$152,400", "$436,350"],
        ["Marketing",     "$89,500",  "$94,200",  "$76,800",  "$260,500"],
        ["Human Resources", "$62,300", "$58,900",  "$64,100",  "$185,300"],
        ["Operations",    "$112,800", "$108,400", "$115,600", "$336,800"],
        ["Sales",         "$78,600",  "$82,300",  "$91,200",  "$252,100"],
        ["Research & Dev", "$134,500", "$142,800", "$128,900", "$406,200"],
    ]

    text_color = RGBColor(0x33, 0x33, 0x33)
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            align = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            set_cell_text(table.cell(r, c), val, font_size=Pt(12),
                          bold=False, color=text_color, alignment=align)
            set_cell_fill(table.cell(r, c), "FFFFFF")

    # --- Slide 4: Analysis ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Trend Analysis"
    body4 = slide4.placeholders[1]
    tf4 = body4.text_frame
    tf4.text = "Engineering maintains the highest spend at $436K, driven by cloud infrastructure expansion and new hires in the platform team."
    p4a = tf4.add_paragraph()
    p4a.text = "Marketing achieved a 19% cost reduction in Q3 by shifting from print to digital campaigns, with improved conversion rates."
    p4b = tf4.add_paragraph()
    p4b.text = "R&D investment remains strong at $406K, supporting three major product initiatives scheduled for Q4 launch."

    # --- Slide 5: Next Steps ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Recommendations & Next Steps"
    body5 = slide5.placeholders[1]
    tf5 = body5.text_frame
    tf5.text = "1. Approve additional $50K for Engineering cloud migration in Q4"
    p5a = tf5.add_paragraph()
    p5a.text = "2. Maintain current Marketing digital-first strategy through year-end"
    p5b = tf5.add_paragraph()
    p5b.text = "3. Review Operations vendor contracts for potential consolidation savings"
    p5c = tf5.add_paragraph()
    p5c.text = "4. Schedule Q4 budget review meeting for November 15, 2025"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
