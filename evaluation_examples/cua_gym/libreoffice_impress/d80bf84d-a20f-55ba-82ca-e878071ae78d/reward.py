"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 76 is driving me nuts—the photo labeled “Picture 1” looks off-balance. In LibreOffice Impress, how do I crop that image to an exact 4:3 aspect ratio and then shift the frame so the main subject sits dead center?
Generated: 2025-09-10 23:59:13
Status: success
Model: azure-o3
Total Steps: 10
"""

import os
from pptx import Presentation


def verify_slide76_crop(file_path: str) -> float:
    """Verify that on slide 76 (1-based index) the picture labelled
    "Picture 1" has been cropped to an exact 4:3 visible aspect ratio and
    that the crop frame is centred on the subject (i.e. symmetrical left /
    right and top / bottom margins).

    Scoring (progressive):
        • 0.1  – target picture found on slide 76
        • 0.2  – any cropping (non-zero srcRect) applied
        • 0.3  – resulting visible area within 3 % tolerance of 4:3
        • 0.4  – crop symmetrical (≤0.5 % difference L/R and T/B) ⇒ centred
        --------------------------------------------
        1.0  – perfect completion
    """

    MAX_SCORE = 1.0
    score = 0.0

    # 1. Open the presentation file (no points – prerequisite)
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Could not open PPTX: {e}")
        return 0.0

    # 2. Ensure slide 76 exists (index 75 because 0-based)
    slide_index = 75
    if slide_index >= len(prs.slides):
        print(f"✗ Slide 76 not found. Presentation has only {len(prs.slides)} slides.")
        return 0.0

    slide = prs.slides[slide_index]

    # 3. Locate the target picture (prefer name startswith "Picture 1")
    picture = None
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            if shape.name.lower().startswith("picture 1"):
                picture = shape
                break
            # Fallback: remember first picture if specific name not found yet
            if picture is None:
                picture = shape

    if picture is None:
        print("✗ No picture found on slide 76")
        return 0.0

    print(f"✓ Located picture: '{picture.name}' (shape id={picture.shape_id})")
    score += 0.1  # picture located

    # 4. Extract cropping information (a:srcRect in p:blipFill)
    blip_fill = picture._element.blipFill
    if blip_fill is None or blip_fill.srcRect is None:
        print("✗ Picture has no cropping applied (srcRect missing)")
        print(f"Total score: {score}")
        return score

    src_rect = blip_fill.srcRect
    attrs = { (k.split('}')[-1] if '}' in k else k): int(v) for k, v in src_rect.attrib.items() }
    l = attrs.get('l', 0)
    r = attrs.get('r', 0)
    t = attrs.get('t', 0)
    b = attrs.get('b', 0)
    print(f"Cropping values (1/1000 %): left={l}, right={r}, top={t}, bottom={b}")

    # 4a. Any crop applied?
    if any([l, r, t, b]):
        print("✓ Cropping detected (values not all zero)")
        score += 0.2
    else:
        print("✗ Cropping values all zero → no crop")
        print(f"Total score: {score}")
        return score

    # 5. Verify visible area aspect ratio ≈ 4:3 (within 3 %)
    visible_w_pct = 100000 - l - r
    visible_h_pct = 100000 - t - b
    if visible_h_pct == 0 or visible_w_pct == 0:
        print("✗ Invalid cropping results in zero dimension")
        print(f"Total score: {score}")
        return score

    # The picture is rendered at picture.width × picture.height EMU, then the
    # srcRect further crops it. Effective ratio = shape_ratio * w%/h%
    shape_ratio = picture.width / picture.height
    effective_ratio = shape_ratio * visible_w_pct / visible_h_pct
    target_ratio = 4 / 3
    diff_ratio = abs(effective_ratio - target_ratio) / target_ratio
    print(f"Visible aspect ratio: {effective_ratio:.3f}  | diff to 4:3 = {diff_ratio:.2%}")

    if diff_ratio <= 0.03:  # within 3 %
        print("✓ Aspect ratio within tolerance of 4:3")
        score += 0.3
    else:
        print("✗ Aspect ratio outside tolerance (>3 %)")

    # 6. Verify centring – symmetrical crop (≤0.5 % difference)
    if abs(l - r) <= 500 and abs(t - b) <= 500:
        print("✓ Cropping symmetrical → subject likely centred")
        score += 0.4
    else:
        print("✗ Cropping not symmetrical enough for centring (differences >0.5 %)")

    final_score = min(score, MAX_SCORE)
    print(f"Total score: {final_score}")
    return final_score


if __name__ == "__main__":
    # Path is fixed in the workspace; adjust here if task location changes.
    pptx_path = "/home/user/slide_76_is_driving_me_nutsthe_photo_labeled_picture_1_looks_off_balance_in_libreoffice_impress_how__golden.pptx"

    reward = verify_slide76_crop(pptx_path)
    print(f"REWARD: {reward}")
