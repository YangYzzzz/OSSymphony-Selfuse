"""
Initial Setup: Create accessible lecture presentation with mixed font sizes
Task ID: impress_teach_043
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_043'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size, font_name="Arial",
                bold=False, color=None, alignment=None):
    """Helper to add a text box with specific font properties."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)
    return txBox


def add_multi_run_textbox(slide, left, top, width, height, runs_data, alignment=None):
    """Add a textbox with multiple runs of different sizes."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if alignment:
        p.alignment = alignment
    for i, (text, size, bold, color) in enumerate(runs_data):
        if i == 0:
            p.text = text
            run = p.runs[0]
        else:
            run = p.add_run()
            run.text = text
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = RGBColor(*color)
    return txBox


def add_bullet_slide(slide, left, top, width, height, items):
    """Add a text box with bullet-like paragraphs, each with its own font size."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, bold) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========== Slide 1: Title Slide ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_textbox(slide1, Inches(1), Inches(1.5), Inches(8), Inches(1.5),
                "Accessible Design in Digital Learning", 32, bold=True,
                color=(0x1A, 0x3C, 0x6D), alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(1.5), Inches(3.5), Inches(7), Inches(1),
                "Department of Educational Technology | Spring 2025", 20,
                color=(0x55, 0x55, 0x55), alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(2), Inches(5), Inches(6), Inches(0.8),
                "Dr. Elena Vasquez, Associate Professor", 22,
                color=(0x33, 0x33, 0x33), alignment=PP_ALIGN.CENTER)

    # ========== Slide 2: Course Overview ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide2, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Course Overview", 28, bold=True,
                color=(0x1A, 0x3C, 0x6D))
    add_bullet_slide(slide2, Inches(0.8), Inches(1.5), Inches(8.4), Inches(5), [
        ("This course explores principles of universal design for learning", 20, False),
        ("We will examine WCAG 2.1 guidelines and their application", 20, False),
        ("Hands-on projects using assistive technologies", 20, False),
        ("Assessment through portfolio-based evaluation", 20, False),
    ])

    # ========== Slide 3: Principles of Accessibility (mixed sizes) ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide3, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Principles of Accessibility", 28, bold=True,
                color=(0x1A, 0x3C, 0x6D))
    # Content with small fonts that need fixing
    add_bullet_slide(slide3, Inches(0.8), Inches(1.5), Inches(8.4), Inches(5), [
        ("Perceivable: Information must be presentable in ways users can perceive", 14, False),
        ("Operable: Interface components must be operable by all users", 12, False),
        ("Understandable: Content must be readable and predictable", 10, False),
        ("Robust: Content must be compatible with assistive technologies", 16, False),
    ])
    # Small footnote
    add_textbox(slide3, Inches(0.8), Inches(6), Inches(8), Inches(0.5),
                "Source: W3C Web Content Accessibility Guidelines 2.1", 10,
                color=(0x88, 0x88, 0x88))

    # ========== Slide 4: Visual Design Guidelines (mixed sizes) ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide4, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Visual Design Guidelines", 24, bold=True,
                color=(0x1A, 0x3C, 0x6D))
    add_bullet_slide(slide4, Inches(0.8), Inches(1.5), Inches(8.4), Inches(4.5), [
        ("Use high contrast ratios (minimum 4.5:1 for normal text)", 16, False),
        ("Avoid relying solely on color to convey meaning", 12, False),
        ("Provide text alternatives for all non-text content", 14, False),
        ("Ensure sufficient spacing between interactive elements", 10, False),
        ("Use scalable vector graphics where possible", 12, False),
    ])
    add_textbox(slide4, Inches(0.8), Inches(6.2), Inches(8), Inches(0.4),
                "Refer to the WCAG color contrast checker tool for verification", 11,
                color=(0x66, 0x66, 0x66))

    # ========== Slide 5: Assistive Technology Overview (mixed sizes) ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide5, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Assistive Technology Overview", 26, bold=True,
                color=(0x1A, 0x3C, 0x6D))
    add_bullet_slide(slide5, Inches(0.8), Inches(1.5), Inches(4), Inches(4), [
        ("Screen Readers", 20, True),
        ("JAWS, NVDA, VoiceOver", 12, False),
        ("Magnification Software", 20, True),
        ("ZoomText, Windows Magnifier", 14, False),
    ])
    add_bullet_slide(slide5, Inches(5.2), Inches(1.5), Inches(4), Inches(4), [
        ("Alternative Input Devices", 20, True),
        ("Switch access, eye tracking, sip-and-puff", 10, False),
        ("Speech Recognition", 20, True),
        ("Dragon NaturallySpeaking, built-in OS tools", 16, False),
    ])

    # ========== Slide 6: Case Study - University Portal Redesign (mixed sizes) ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide6, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Case Study: University Portal Redesign", 24, bold=True,
                color=(0x1A, 0x3C, 0x6D))
    add_multi_run_textbox(slide6, Inches(0.8), Inches(1.5), Inches(8.4), Inches(1.5), [
        ("In 2023, Riverside University redesigned their student portal after an accessibility audit revealed ", 14, False, None),
        ("47 critical violations", 14, True, (0xCC, 0x00, 0x00)),
        (" of WCAG 2.1 Level AA standards.", 14, False, None),
    ])
    add_bullet_slide(slide6, Inches(0.8), Inches(3.2), Inches(8.4), Inches(3.5), [
        ("Navigation restructured with ARIA landmarks and skip links", 12, False),
        ("All images received descriptive alt text reviewed by disability services", 10, False),
        ("Form validation messages made accessible to screen readers", 16, False),
        ("Color palette updated to meet 4.5:1 contrast minimum", 14, False),
        ("Result: 92% compliance score, up from 34%", 20, True),
    ])

    # ========== Slide 7: Document Accessibility Checklist (mixed sizes) ==========
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide7, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Document Accessibility Checklist", 26, bold=True,
                color=(0x1A, 0x3C, 0x6D))
    add_bullet_slide(slide7, Inches(0.8), Inches(1.5), Inches(8.4), Inches(4.5), [
        ("Use heading styles (H1-H6) instead of manually formatted bold text", 14, False),
        ("Include table headers and captions for data tables", 12, False),
        ("Add meaningful hyperlink text (avoid 'click here')", 10, False),
        ("Ensure reading order matches visual layout", 16, False),
        ("Test with at least one screen reader before publishing", 12, False),
        ("Provide document language metadata", 10, False),
    ])
    add_textbox(slide7, Inches(0.8), Inches(6.2), Inches(8), Inches(0.4),
                "Complete checklist available on the course LMS under Resources", 11,
                color=(0x66, 0x66, 0x66))

    # ========== Slide 8: Legal Framework and Standards (mixed sizes) ==========
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide8, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Legal Framework and Standards", 28, bold=True,
                color=(0x1A, 0x3C, 0x6D))
    add_bullet_slide(slide8, Inches(0.8), Inches(1.5), Inches(8.4), Inches(4.5), [
        ("Section 508 of the Rehabilitation Act (United States)", 16, False),
        ("Americans with Disabilities Act (ADA) Title II and III", 14, False),
        ("EN 301 549 (European accessibility standard)", 12, False),
        ("AODA - Accessibility for Ontarians with Disabilities Act", 10, False),
        ("UN Convention on the Rights of Persons with Disabilities", 14, False),
    ])
    add_textbox(slide8, Inches(0.8), Inches(6), Inches(8), Inches(0.5),
                "Note: Non-compliance penalties can exceed $75,000 per violation", 12,
                color=(0x99, 0x33, 0x33))

    # ========== Slide 9: Assignment Instructions ==========
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide9, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Assignment: Accessibility Audit", 28, bold=True,
                color=(0x1A, 0x3C, 0x6D))
    add_bullet_slide(slide9, Inches(0.8), Inches(1.5), Inches(8.4), Inches(5), [
        ("Select a public website and perform a full WCAG 2.1 Level AA audit", 20, False),
        ("Document at least 10 accessibility issues with screenshots", 20, False),
        ("Propose remediation strategies for each identified issue", 20, False),
        ("Submit a 5-page report with executive summary by March 28", 20, False),
    ])

    # ========== Slide 10: Contact and Resources ==========
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide10, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Contact and Resources", 28, bold=True,
                color=(0x1A, 0x3C, 0x6D))
    add_bullet_slide(slide10, Inches(0.8), Inches(1.5), Inches(8.4), Inches(5), [
        ("Office Hours: Tuesday and Thursday, 2:00-4:00 PM, Room 312B", 20, False),
        ("Email: e.vasquez@riverside.edu", 20, False),
        ("W3C WCAG 2.1: https://www.w3.org/TR/WCAG21/", 18, False),
        ("WebAIM: https://webaim.org/", 18, False),
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
