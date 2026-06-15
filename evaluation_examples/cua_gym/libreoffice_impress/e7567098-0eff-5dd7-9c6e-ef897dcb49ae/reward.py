"""
Reward Script: Set all text on slides 3-8 to minimum 18pt font size
Task ID: impress_teach_043
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): All text on slides 3-8 is >= 18pt
  Component 2 (0.3): Previously small runs are now exactly 18pt (not over-enlarged)
  Component 3 (0.3): Previously large runs (>= 18pt) are preserved unchanged
"""

import os

from pptx import Presentation
from pptx.util import Pt

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_043'

# Known initial font sizes for slides 3-8 runs (0-indexed slides 2-7)
# Collected from initial_env exploration. Each entry: (slide_0idx, shape_id, run_index_in_shape, initial_pt)
# We use these to verify that large text wasn't changed and small text was set to exactly 18pt.
# Instead of hardcoding every run, we categorize by checking the threshold.

# Initial sizes of ALL runs on slides 3-8 (from exploration), keyed by (slide_1based, shape_id, run_text_prefix)
# We'll use a simpler approach: count-based verification.

# From initial_env:
# Slides 3-8 have 35 runs < 18pt and 11 runs >= 18pt
INITIAL_SMALL_COUNT = 35
INITIAL_LARGE_COUNT = 11

def get_slide_runs(prs, slide_indices):
    """Get all non-empty text runs from specified slides (0-indexed)."""
    runs_info = []
    for idx in slide_indices:
        slide = prs.slides[idx]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            sz = run.font.size
                            pt = round(sz / 12700, 1) if sz else 0
                            runs_info.append({
                                'slide': idx + 1,
                                'shape_id': shape.shape_id,
                                'text': run.text[:40],
                                'size_pt': pt,
                            })
    return runs_info


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Verify slide count is still 10
    if len(prs.slides) != 10:
        print(f"CRITICAL: Expected 10 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Get all runs on slides 3-8 (0-indexed: 2-7)
    target_slides = list(range(2, 8))
    runs = get_slide_runs(prs, target_slides)

    if not runs:
        print("CRITICAL: No text runs found on slides 3-8")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Found {len(runs)} text runs on slides 3-8")

    # Component 1: All text on slides 3-8 is >= 18pt (0.4 points)
    # This is the PRIMARY requirement. In initial_env, 35 runs are < 18pt, so this FAILS.
    try:
        below_18 = [r for r in runs if r['size_pt'] < 18.0]
        if len(below_18) == 0:
            print(f"PASS: Component 1 - All {len(runs)} runs on slides 3-8 are >= 18pt (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 - {len(below_18)} runs still below 18pt:")
            for r in below_18[:5]:
                print(f"  Slide {r['slide']}, shape {r['shape_id']}: '{r['text']}' = {r['size_pt']}pt")
            if len(below_18) > 5:
                print(f"  ... and {len(below_18) - 5} more")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Runs that should be exactly 18pt (were < 18pt initially) are 18pt (0.3 points)
    # In the golden, all originally-small runs should be exactly 18.0pt, not some other value.
    # We check: count of runs at exactly 18pt should match the initial small count (35).
    try:
        at_18 = [r for r in runs if abs(r['size_pt'] - 18.0) < 0.1]
        # Also need runs > 18pt that were already large
        above_18 = [r for r in runs if r['size_pt'] > 18.1]

        # In initial: 35 small + 11 large = 46 total runs
        # In golden: should be 35 at 18pt + 11 at original large sizes = 46 total
        # Score based on how many of the expected 35 are at exactly 18pt
        if INITIAL_SMALL_COUNT > 0:
            ratio = min(len(at_18), INITIAL_SMALL_COUNT) / INITIAL_SMALL_COUNT
            # Only award if ratio is high (most small runs are now 18pt)
            if ratio >= 0.9:
                print(f"PASS: Component 2 - {len(at_18)} runs at exactly 18pt (expected ~{INITIAL_SMALL_COUNT}) (0.3 pts)")
                total_score += 0.3
            elif ratio >= 0.5:
                partial = round(0.3 * ratio, 2)
                print(f"PARTIAL: Component 2 - {len(at_18)} runs at 18pt vs expected {INITIAL_SMALL_COUNT}, ratio={ratio:.2f} ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 2 - Only {len(at_18)} runs at 18pt, expected ~{INITIAL_SMALL_COUNT}")
        else:
            print(f"SKIP: Component 2 - No initially small runs expected")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Text already >= 18pt is preserved AND small text was raised (0.3 points)
    # This is a compound check: it only awards points if the minimum-size requirement (Component 1)
    # is also met, ensuring this component FAILS on initial_env where small runs still exist.
    try:
        if len(below_18) > 0:
            # If there are still runs below 18pt, this compound check fails
            print(f"FAIL: Component 3 - Cannot verify preservation when {len(below_18)} runs are still < 18pt")
        else:
            # All runs >= 18pt. Now check that the originally-large runs kept their sizes.
            initial_large_sizes = {28.0, 24.0, 26.0, 20.0}
            large_runs = [r for r in runs if r['size_pt'] > 18.1]

            if len(large_runs) >= INITIAL_LARGE_COUNT - 1:
                all_valid = all(r['size_pt'] in initial_large_sizes for r in large_runs)
                if all_valid:
                    print(f"PASS: Component 3 - {len(large_runs)} large runs preserved at original sizes (0.3 pts)")
                    total_score += 0.3
                else:
                    partial = 0.15
                    print(f"PARTIAL: Component 3 - Large runs exist but some sizes differ ({partial} pts)")
                    total_score += partial
            else:
                print(f"FAIL: Component 3 - Only {len(large_runs)} large runs (>18pt), expected ~{INITIAL_LARGE_COUNT}")
                print(f"  This suggests large text was incorrectly changed to 18pt")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
