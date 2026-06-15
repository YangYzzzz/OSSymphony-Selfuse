"""
Reward Script: Create a styled chemistry elements table on slide 3
Task ID: impress_stu_033
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Table exists on slide 3 with 6 rows x 5 columns
  Component 2 (0.35): Table data rows contain correct element data
  Component 3 (0.20): Header row cells have dark blue (#003366) background
  Component 4 (0.20): Header row text is white (#FFFFFF)
"""

import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_033'

# Expected table data (header + 5 data rows)
EXPECTED_HEADERS = ['Element', 'Symbol', 'Atomic Number', 'Electronegativity', 'Category']
EXPECTED_DATA = [
    ['Hydrogen', 'H', '1', '2.20', 'Nonmetal'],
    ['Carbon', 'C', '6', '2.55', 'Nonmetal'],
    ['Nitrogen', 'N', '7', '3.04', 'Nonmetal'],
    ['Oxygen', 'O', '8', '3.44', 'Nonmetal'],
    ['Iron', 'Fe', '26', '1.83', 'Metal'],
]


def find_table_on_slide(slide):
    """Find first table shape on a slide, return table object or None."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return shape.table
    return None


def get_cell_bg_color(cell):
    """Extract cell background color as uppercase hex string (e.g. '003366'), or None."""
    try:
        tc_elem = cell._tc
        tcPr = tc_elem.find(qn('a:tcPr'))
        if tcPr is not None:
            solidFill = tcPr.find(qn('a:solidFill'))
            if solidFill is not None:
                srgb = solidFill.find(qn('a:srgbClr'))
                if srgb is not None:
                    return srgb.get('val').upper()
    except Exception:
        pass
    return None


def get_cell_text_color(cell):
    """Extract text color from first run in cell as uppercase hex string, or None."""
    try:
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                if run.font.color.type is not None:
                    return str(run.font.color.rgb).upper()
    except Exception:
        pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]  # 0-indexed
    table = find_table_on_slide(slide3)

    # Component 1: Table exists on slide 3 with correct dimensions (0.25 points)
    try:
        if table is None:
            print("FAIL: Component 1 — No table found on slide 3")
        else:
            nrows = len(table.rows)
            ncols = len(table.columns)
            if nrows == 6 and ncols == 5:
                print(f"PASS: Component 1 — Table found on slide 3 with {nrows} rows x {ncols} cols (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 — Table dimensions {nrows}x{ncols}, expected 6x5")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # If no table, remaining checks are impossible
    if table is None:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: Table data correctness (0.35 points)
    # Check header row + 5 data rows. Award partial credit per correct row.
    try:
        data_correct = 0
        total_rows_to_check = 6  # 1 header + 5 data

        # Check header row
        header_vals = [table.cell(0, c).text.strip() for c in range(min(len(table.columns), 5))]
        if header_vals == EXPECTED_HEADERS:
            data_correct += 1
            print(f"PASS: Component 2a — Header row correct: {header_vals}")
        else:
            print(f"FAIL: Component 2a — Header row: {header_vals}, expected: {EXPECTED_HEADERS}")

        # Check data rows
        for row_idx, expected_row in enumerate(EXPECTED_DATA):
            if row_idx + 1 >= len(table.rows):
                print(f"FAIL: Component 2b — Row {row_idx + 1} missing (not enough rows)")
                continue
            actual_row = [table.cell(row_idx + 1, c).text.strip() for c in range(min(len(table.columns), 5))]
            if actual_row == expected_row:
                data_correct += 1
                print(f"PASS: Component 2b — Row {row_idx + 1} correct: {actual_row}")
            else:
                print(f"FAIL: Component 2b — Row {row_idx + 1}: {actual_row}, expected: {expected_row}")

        row_score = 0.35 * (data_correct / total_rows_to_check)
        if row_score > 0:
            print(f"PASS: Component 2 — {data_correct}/{total_rows_to_check} rows correct ({row_score:.3f} pts)")
            total_score += row_score
        else:
            print(f"FAIL: Component 2 — No rows matched")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Header row background color is dark blue #003366 (0.20 points)
    try:
        bg_correct = 0
        for c in range(min(len(table.columns), 5)):
            bg = get_cell_bg_color(table.cell(0, c))
            if bg == '003366':
                bg_correct += 1
            else:
                print(f"FAIL: Component 3 — Cell (0,{c}) bg color: {bg}, expected: 003366")

        if bg_correct == 5:
            print(f"PASS: Component 3 — All 5 header cells have #003366 background (0.20 pts)")
            total_score += 0.20
        elif bg_correct > 0:
            partial = 0.20 * (bg_correct / 5)
            print(f"PARTIAL: Component 3 — {bg_correct}/5 header cells have correct bg ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No header cells have #003366 background")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Header row text is white #FFFFFF (0.20 points)
    try:
        txt_correct = 0
        for c in range(min(len(table.columns), 5)):
            txt_color = get_cell_text_color(table.cell(0, c))
            if txt_color == 'FFFFFF':
                txt_correct += 1
            else:
                print(f"FAIL: Component 4 — Cell (0,{c}) text color: {txt_color}, expected: FFFFFF")

        if txt_correct == 5:
            print(f"PASS: Component 4 — All 5 header cells have white text (0.20 pts)")
            total_score += 0.20
        elif txt_correct > 0:
            partial = 0.20 * (txt_correct / 5)
            print(f"PARTIAL: Component 4 — {txt_correct}/5 header cells have white text ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — No header cells have white text")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
