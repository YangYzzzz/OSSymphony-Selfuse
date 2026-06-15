"""
Initial Setup: Chemistry Review presentation with 8 slides, slide 3 empty except for title
Task ID: impress_stu_033
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_033'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_content(slide, title_text, body_lines):
    """Helper to set title and body content on a Title+Content layout slide."""
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Chemistry Review"
    slide1.placeholders[1].text = "General Chemistry - Spring 2025"

    # --- Slide 2: Atomic Structure ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_content(slide2, "Atomic Structure", [
        "Atoms consist of protons, neutrons, and electrons",
        "Protons carry positive charge (+1)",
        "Neutrons are electrically neutral",
        "Electrons carry negative charge (-1)",
        "Atomic number = number of protons",
        "Mass number = protons + neutrons",
    ])

    # --- Slide 3: Key Elements (EMPTY - only title) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Elements"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 4: The Periodic Table ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_content(slide4, "The Periodic Table", [
        "Organized by atomic number",
        "Groups (columns) share similar chemical properties",
        "Periods (rows) show trends in atomic radius",
        "Metals on the left, nonmetals on the right",
        "Transition metals in the center block",
    ])

    # --- Slide 5: Chemical Bonding ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_content(slide5, "Chemical Bonding", [
        "Ionic bonds: transfer of electrons between atoms",
        "Covalent bonds: sharing of electron pairs",
        "Metallic bonds: delocalized electron sea model",
        "Bond polarity depends on electronegativity difference",
    ])

    # --- Slide 6: Electronegativity ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_content(slide6, "Electronegativity", [
        "Measure of an atom's ability to attract bonding electrons",
        "Pauling scale ranges from 0.7 (Francium) to 3.98 (Fluorine)",
        "Increases across a period (left to right)",
        "Decreases down a group (top to bottom)",
        "Important for predicting bond type and polarity",
    ])

    # --- Slide 7: Chemical Reactions ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_content(slide7, "Chemical Reactions", [
        "Synthesis: A + B -> AB",
        "Decomposition: AB -> A + B",
        "Single replacement: A + BC -> AC + B",
        "Double replacement: AB + CD -> AD + CB",
        "Combustion: fuel + O2 -> CO2 + H2O",
    ])

    # --- Slide 8: Summary ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_content(slide8, "Summary & Next Steps", [
        "Review atomic structure fundamentals",
        "Practice identifying elements on the periodic table",
        "Understand bonding types and electronegativity trends",
        "Next topic: Stoichiometry and molar calculations",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
