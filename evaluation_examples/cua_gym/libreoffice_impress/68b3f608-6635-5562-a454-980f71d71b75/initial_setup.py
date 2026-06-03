"""
Initial Setup: Export presentation as PDF
Task ID: impress_sales_040
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_040'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tf


def add_bullet_list(tf, items, font_size=16, color=None):
    """Add bullet items to an existing text frame."""
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        run = p.runs[0]
        run.font.size = Pt(font_size)
        if color:
            run.font.color.rgb = color


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    DARK_BG = RGBColor(0x1B, 0x2A, 0x4A)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    ACCENT_BLUE = RGBColor(0x00, 0x9E, 0xDB)
    LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
    DARK_TEXT = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    add_text_box(slide1, Inches(1), Inches(1.5), Inches(11), Inches(2),
                 "CloudSync", font_size=54, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, Inches(1), Inches(3.5), Inches(11), Inches(1),
                 "Enterprise Cloud Synchronization Platform", font_size=28,
                 color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, Inches(1), Inches(5.0), Inches(11), Inches(1),
                 "Proposal for Meridian Financial Group  |  Q2 2026",
                 font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide2, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Agenda", font_size=36, bold=True, color=DARK_BG)
    tf2 = add_text_box(slide2, Inches(0.8), Inches(1.6), Inches(11), Inches(5),
                       "1. The Challenge of Data Fragmentation", font_size=20,
                       color=DARK_TEXT)
    add_bullet_list(tf2, [
        "2. CloudSync Platform Overview",
        "3. Key Features & Differentiators",
        "4. Security & Compliance Framework",
        "5. Integration Ecosystem",
        "6. Pricing & Licensing",
        "7. Implementation Timeline",
        "8. Client Success Stories",
        "9. Next Steps & Q&A"
    ], font_size=20, color=DARK_TEXT)

    # --- Slide 3: The Problem ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide3, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "The Challenge", font_size=36, bold=True, color=DARK_BG)
    tf3 = add_text_box(slide3, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5),
                       "73% of enterprises struggle with data silos across cloud providers.",
                       font_size=18, color=DARK_TEXT)
    add_bullet_list(tf3, [
        "Average company uses 4.2 cloud providers simultaneously",
        "Data inconsistency costs $12.9M annually per Fortune 500 firm",
        "IT teams spend 38% of time on manual sync and reconciliation",
        "Security gaps from fragmented access control policies"
    ], font_size=16, color=DARK_TEXT)

    add_text_box(slide3, Inches(7.0), Inches(1.6), Inches(5.5), Inches(5),
                 "Industry Pain Points:\n\n"
                 "- File version conflicts across teams\n"
                 "- Compliance audit failures due to scattered data\n"
                 "- Slow disaster recovery from fragmented backups\n"
                 "- Shadow IT proliferation from poor sync tools",
                 font_size=16, color=DARK_TEXT)

    # --- Slide 4: Solution Overview ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    fill4 = slide4.background.fill
    fill4.solid()
    fill4.fore_color.rgb = DARK_BG
    add_text_box(slide4, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "CloudSync Platform Overview", font_size=36, bold=True, color=WHITE)
    add_text_box(slide4, Inches(0.8), Inches(1.8), Inches(11), Inches(1.5),
                 "A unified synchronization layer that connects AWS, Azure, Google Cloud, "
                 "and on-premise storage into a single, real-time data fabric.",
                 font_size=20, color=ACCENT_BLUE)
    tf4 = add_text_box(slide4, Inches(0.8), Inches(3.5), Inches(5.5), Inches(3.5),
                       "Core Architecture:", font_size=22, bold=True, color=WHITE)
    add_bullet_list(tf4, [
        "Distributed sync engine with <50ms latency",
        "Conflict resolution via vector clock algorithm",
        "End-to-end AES-256 encryption in transit and at rest",
        "Zero-trust access model with SAML 2.0 / OAuth"
    ], font_size=16, color=LIGHT_GRAY)

    tf4b = add_text_box(slide4, Inches(7.0), Inches(3.5), Inches(5.5), Inches(3.5),
                        "Deployment Options:", font_size=22, bold=True, color=WHITE)
    add_bullet_list(tf4b, [
        "SaaS (fully managed)",
        "Hybrid (control plane in your VPC)",
        "On-premise (air-gapped environments)",
        "Kubernetes-native orchestration"
    ], font_size=16, color=LIGHT_GRAY)

    # --- Slide 5: Key Features ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide5, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Key Features & Differentiators", font_size=36, bold=True, color=DARK_BG)

    features = [
        ("Real-Time Sync", "Sub-second synchronization across up to 12 cloud providers "
         "with automatic conflict detection and resolution."),
        ("Smart Deduplication", "AI-powered content-aware deduplication reduces storage "
         "costs by up to 60% while maintaining full file fidelity."),
        ("Compliance Dashboard", "Pre-built templates for SOC 2, HIPAA, GDPR, and "
         "ISO 27001 with automated audit trail generation."),
    ]
    for i, (title, desc) in enumerate(features):
        y = Inches(1.6 + i * 1.8)
        add_text_box(slide5, Inches(0.8), y, Inches(3), Inches(0.6),
                     title, font_size=22, bold=True, color=ACCENT_BLUE)
        add_text_box(slide5, Inches(0.8), y + Inches(0.6), Inches(11), Inches(1),
                     desc, font_size=16, color=DARK_TEXT)

    # --- Slide 6: Security & Compliance ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide6, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Security & Compliance", font_size=36, bold=True, color=DARK_BG)
    tf6 = add_text_box(slide6, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5),
                       "Certifications & Standards:", font_size=22, bold=True,
                       color=DARK_TEXT)
    add_bullet_list(tf6, [
        "SOC 2 Type II certified",
        "HIPAA BAA available",
        "GDPR compliant with EU data residency",
        "FedRAMP Moderate (in progress)",
        "ISO 27001:2022 certified"
    ], font_size=16, color=DARK_TEXT)

    tf6b = add_text_box(slide6, Inches(7.0), Inches(1.6), Inches(5.5), Inches(5),
                        "Security Features:", font_size=22, bold=True, color=DARK_TEXT)
    add_bullet_list(tf6b, [
        "Zero-knowledge encryption option",
        "Hardware security module (HSM) key management",
        "Role-based access control with 47 granular permissions",
        "Real-time threat detection and anomaly alerts",
        "Immutable audit logs with 7-year retention"
    ], font_size=16, color=DARK_TEXT)

    # --- Slide 7: Pricing ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide7, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Pricing & Licensing", font_size=36, bold=True, color=DARK_BG)

    # Create pricing table
    table_shape = slide7.shapes.add_table(5, 4, Inches(0.8), Inches(1.6),
                                          Inches(11.5), Inches(4.5))
    table = table_shape.table
    headers = ["Plan", "Users", "Storage", "Annual Price"]
    data_rows = [
        ["Starter", "Up to 50", "500 GB", "$18,000"],
        ["Professional", "Up to 250", "5 TB", "$72,000"],
        ["Enterprise", "Unlimited", "Unlimited", "$156,000"],
        ["Enterprise Plus", "Unlimited", "Unlimited + DR", "Custom"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(16)
            run.font.color.rgb = WHITE
    for r, row in enumerate(data_rows, 1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(14)

    # --- Slide 8: Implementation Timeline ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide8, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Implementation Timeline", font_size=36, bold=True, color=DARK_BG)
    phases = [
        ("Week 1-2", "Discovery & Planning",
         "Requirements gathering, architecture review, and integration mapping"),
        ("Week 3-4", "Core Setup",
         "Platform deployment, SSO configuration, and initial data migration"),
        ("Week 5-6", "Integration & Testing",
         "API integrations, UAT testing, and performance benchmarking"),
        ("Week 7-8", "Go-Live & Optimization",
         "Production cutover, monitoring setup, and team training sessions"),
    ]
    for i, (period, title, desc) in enumerate(phases):
        y = Inches(1.6 + i * 1.4)
        add_text_box(slide8, Inches(0.8), y, Inches(2), Inches(0.5),
                     period, font_size=18, bold=True, color=ACCENT_BLUE)
        add_text_box(slide8, Inches(3.0), y, Inches(3.5), Inches(0.5),
                     title, font_size=18, bold=True, color=DARK_TEXT)
        add_text_box(slide8, Inches(3.0), y + Inches(0.5), Inches(9), Inches(0.7),
                     desc, font_size=14, color=DARK_TEXT)

    # --- Slide 9: Client Testimonials ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    fill9 = slide9.background.fill
    fill9.solid()
    fill9.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
    add_text_box(slide9, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "What Our Clients Say", font_size=36, bold=True, color=DARK_BG)

    testimonials = [
        ('"CloudSync reduced our cross-cloud sync errors by 94% in the first quarter."',
         "- Rachel Torres, VP of Engineering, Pinnacle Health Systems"),
        ('"The compliance dashboard alone saved us 200+ hours during our SOC 2 audit."',
         "- David Nakamura, CISO, Westfield Capital Partners"),
        ('"We consolidated 6 different sync tools into one platform. ROI was evident in 3 months."',
         "- Samantha Okafor, CTO, NovaTech Industries"),
    ]
    for i, (quote, author) in enumerate(testimonials):
        y = Inches(1.6 + i * 1.8)
        add_text_box(slide9, Inches(1.2), y, Inches(10.5), Inches(1),
                     quote, font_size=18, color=DARK_TEXT)
        add_text_box(slide9, Inches(1.2), y + Inches(1.0), Inches(10.5), Inches(0.5),
                     author, font_size=14, bold=True, color=ACCENT_BLUE)

    # --- Slide 10: Next Steps / Closing ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    fill10 = slide10.background.fill
    fill10.solid()
    fill10.fore_color.rgb = DARK_BG
    add_text_box(slide10, Inches(1), Inches(1.0), Inches(11), Inches(1.5),
                 "Next Steps", font_size=44, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    tf10 = add_text_box(slide10, Inches(1.5), Inches(2.8), Inches(10), Inches(3),
                        "1. Schedule a technical deep-dive with your IT team",
                        font_size=22, color=WHITE)
    add_bullet_list(tf10, [
        "2. Pilot program with 25 users for 30 days (no cost)",
        "3. Integration assessment with your existing stack",
        "4. Executive review and contract discussion"
    ], font_size=22, color=WHITE)
    add_text_box(slide10, Inches(1), Inches(5.8), Inches(11), Inches(1),
                 "Contact: Alexandra Voss  |  avoss@cloudsync.io  |  +1 (415) 555-0192",
                 font_size=18, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
