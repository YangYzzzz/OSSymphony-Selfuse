"""
Reward Script: Delete slides 8, 9, and 10 from presentation
Task ID: impress_teach_010
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Slide count is exactly 9
  Component 2 (0.3): Deleted slides' content (Oligopoly, Game Theory, Pure Monopoly) is absent
  Component 3 (0.3): Original slides 11-12 are now slides 8-9 with correct titles
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_010'


def persist_app_state(domain):
    """Save any unsaved GUI edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            import time
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_slide_first_text(slide):
    """Get the first non-empty text from a slide's shapes."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                return text
    return ""


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    print(f"INFO: Presentation has {num_slides} slides")

    # Collect all slide texts for analysis
    slide_texts = []
    for i, slide in enumerate(prs.slides):
        text = get_slide_first_text(slide)
        slide_texts.append(text)
        print(f"  Slide {i+1}: {text[:60]}")

    # Component 1: Slide count is exactly 9 (0.4 points)
    # Initial has 12 slides; after deleting 3, should have 9
    try:
        if num_slides == 9:
            print(f"PASS: Component 1 -- Slide count is 9 (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 -- Expected 9 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Deleted slides' content is absent (0.3 points)
    # Slides 8-10 had titles containing "Oligopoly", "Game Theory", "Pure Monopoly"
    try:
        deleted_keywords = ["Oligopoly", "Game Theory", "Pure Monopoly"]
        all_text_combined = " ".join(slide_texts)
        found_deleted = []
        for kw in deleted_keywords:
            if kw.lower() in all_text_combined.lower():
                found_deleted.append(kw)

        if len(found_deleted) == 0:
            print(f"PASS: Component 2 -- Deleted slide content absent (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 -- Found deleted content still present: {found_deleted}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Original slides 11-12 are now slides 8-9 (0.3 points)
    # Slide 8 should be "Practice Problems", Slide 9 should be "References and Further Reading"
    try:
        checks_passed = 0

        if num_slides >= 8:
            slide8_text = slide_texts[7]  # 0-indexed
            if "Practice Problems" in slide8_text:
                checks_passed += 1
                print(f"  Slide 8 text OK: '{slide8_text[:50]}'")
            else:
                print(f"  Slide 8 text MISMATCH: expected 'Practice Problems', found '{slide8_text[:50]}'")
        else:
            print(f"  Not enough slides to check slide 8")

        if num_slides >= 9:
            slide9_text = slide_texts[8]  # 0-indexed
            if "References" in slide9_text:
                checks_passed += 1
                print(f"  Slide 9 text OK: '{slide9_text[:50]}'")
            else:
                print(f"  Slide 9 text MISMATCH: expected 'References...', found '{slide9_text[:50]}'")
        else:
            print(f"  Not enough slides to check slide 9")

        if checks_passed == 2:
            print(f"PASS: Component 3 -- Slides 8-9 have correct content (0.3 pts)")
            total_score += 0.3
        elif checks_passed == 1:
            print(f"PARTIAL: Component 3 -- 1/2 slide positions correct (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 3 -- Slide 8-9 content does not match expected")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
