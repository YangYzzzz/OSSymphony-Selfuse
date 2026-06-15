"""
Initial Setup: Create a 12-slide economics lecture presentation
Task ID: impress_teach_010
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_010'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.paragraphs[0].text = point
        else:
            p = tf.add_paragraph()
            p.text = point
    return slide


def create_initial():
    prs = Presentation()

    # Slide 1: Title Slide
    add_title_slide(
        prs,
        "Economics 101: Lecture 5",
        "Market Structures and Competition\nProfessor Sarah Mitchell\nSpring 2026"
    )

    # Slide 2: Lecture Overview
    add_content_slide(prs, "Lecture Overview", [
        "Review of supply and demand fundamentals",
        "Perfect competition characteristics",
        "Monopolistic competition analysis",
        "Real-world market examples",
        "Practice problems and applications",
    ])

    # Slide 3: Supply and Demand Review
    add_content_slide(prs, "Supply and Demand Review", [
        "Law of demand: price increases lead to quantity decreases",
        "Law of supply: price increases lead to quantity increases",
        "Equilibrium occurs where supply meets demand",
        "Shifts vs. movements along the curve",
        "Consumer and producer surplus at equilibrium",
    ])

    # Slide 4: Perfect Competition - Definition
    add_content_slide(prs, "Perfect Competition: Key Features", [
        "Many buyers and sellers in the market",
        "Homogeneous (identical) products across firms",
        "Free entry and exit from the market",
        "Perfect information available to all participants",
        "Price takers: no single firm influences market price",
    ])

    # Slide 5: Perfect Competition - Examples
    add_content_slide(prs, "Perfect Competition: Real-World Examples", [
        "Agricultural commodities (wheat, corn, soybeans)",
        "Foreign exchange markets for major currencies",
        "Online retail marketplaces with standardized goods",
        "Note: true perfect competition is rare in practice",
        "Most examples are approximations of the model",
    ])

    # Slide 6: Monopolistic Competition - Definition
    add_content_slide(prs, "Monopolistic Competition: Key Features", [
        "Many firms competing in the market",
        "Differentiated products (branding, quality, features)",
        "Low barriers to entry and exit",
        "Some price-setting power due to product differences",
        "Non-price competition through advertising and innovation",
    ])

    # Slide 7: Monopolistic Competition - Analysis
    add_content_slide(prs, "Monopolistic Competition: Short-Run Analysis", [
        "Firms face downward-sloping demand curves",
        "Profit maximization at MR = MC",
        "Short-run economic profits attract new entrants",
        "Long-run equilibrium: zero economic profit",
        "Excess capacity compared to perfect competition",
    ])

    # Slide 8: Oligopoly Introduction (NEXT WEEK)
    add_content_slide(prs, "Oligopoly Markets: Introduction", [
        "Few large firms dominate the market",
        "Significant barriers to entry exist",
        "Strategic interdependence between firms",
        "Products may be homogeneous or differentiated",
        "Examples: automobile, airline, and tech industries",
    ])

    # Slide 9: Game Theory Basics (NEXT WEEK)
    add_content_slide(prs, "Game Theory and Strategic Behavior", [
        "Nash equilibrium: no player benefits from changing strategy",
        "Prisoner's dilemma applied to pricing decisions",
        "Dominant strategies and best response functions",
        "Collusion vs. competition in oligopoly markets",
        "Cartels and their inherent instability",
    ])

    # Slide 10: Monopoly Overview (NEXT WEEK)
    add_content_slide(prs, "Pure Monopoly: Preview for Next Week", [
        "Single seller controls entire market supply",
        "Unique product with no close substitutes",
        "High barriers to entry protect monopolist position",
        "Price maker: firm sets its own price",
        "Deadweight loss and market inefficiency",
    ])

    # Slide 11: Practice Problems
    add_content_slide(prs, "Practice Problems", [
        "1. Identify the market structure for each scenario provided",
        "2. Calculate equilibrium price and quantity given supply and demand equations",
        "3. Compare long-run outcomes in perfect vs. monopolistic competition",
        "4. Explain why restaurants are considered monopolistically competitive",
        "5. Problem set #5 due next Tuesday by 11:59 PM",
    ])

    # Slide 12: References and Reading
    add_content_slide(prs, "References and Further Reading", [
        "Mankiw, N.G. (2024). Principles of Economics, 10th Edition, Ch. 14-16",
        "Varian, H. (2023). Intermediate Microeconomics, Ch. 24-25",
        "The Economist: 'Competition in the Digital Age' (March 2026)",
        "Office hours: Tuesday and Thursday, 2:00-4:00 PM, Room 312",
        "Next lecture: Oligopoly, Game Theory, and Monopoly",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
