"""
Initial Setup: Create Design_Portfolio.pptx with 8 slides, slide 3 titled 'Recent Work' on dark gray background
Task ID: impress_ps_026
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_026'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                font_color=RGBColor(0xFF, 0xFF, 0xFF), bold=False,
                alignment=PP_ALIGN.LEFT, font_name="Arial"):
    """Helper to add a styled text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    return txBox


def set_dark_bg(slide, r=0x2D, g=0x2D, b=0x2D):
    """Set a dark gray background on a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_dark_bg(slide1, 0x1A, 0x1A, 0x2E)
    add_textbox(slide1, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
                "Design Portfolio", font_size=44, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(2.5), Inches(4.0), Inches(8), Inches(1.0),
                "Elena Vasquez — Visual & UX Designer", font_size=20,
                font_color=RGBColor(0xBB, 0xBB, 0xBB),
                alignment=PP_ALIGN.CENTER)

    # --- Slide 2: About Me ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_bg(slide2, 0x23, 0x23, 0x23)
    add_textbox(slide2, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8),
                "About Me", font_size=32, bold=True)
    add_textbox(slide2, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5),
                "I am a multidisciplinary designer with over 8 years of experience "
                "in branding, UI/UX, and motion graphics. My work spans clients "
                "from early-stage startups to Fortune 500 companies including "
                "Rivian, Stripe, and Figma. I believe great design tells a story "
                "and solves real problems.",
                font_size=14, font_color=RGBColor(0xCC, 0xCC, 0xCC))
    add_textbox(slide2, Inches(7.0), Inches(1.8), Inches(5), Inches(4.0),
                "Skills\n• Brand Identity\n• UI/UX Design\n• Motion Graphics\n"
                "• Prototyping (Figma, Sketch)\n• Design Systems\n• Illustration",
                font_size=14, font_color=RGBColor(0xCC, 0xCC, 0xCC))

    # --- Slide 3: Recent Work (TARGET SLIDE — title only, NO body content) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_bg(slide3, 0x2D, 0x2D, 0x2D)
    add_textbox(slide3, Inches(0.8), Inches(0.4), Inches(6), Inches(0.8),
                "Recent Work", font_size=32, bold=True)

    # --- Slide 4: Rivian Rebrand ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_bg(slide4, 0x1E, 0x1E, 0x1E)
    add_textbox(slide4, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
                "Rivian — Brand Refresh", font_size=28, bold=True)
    add_textbox(slide4, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.5),
                "Led the visual identity overhaul for Rivian's 2025 product launch. "
                "Deliverables included a refreshed logo mark, updated color palette "
                "inspired by Pacific Northwest landscapes, a custom typeface "
                "(Rivian Grotesk), and a comprehensive brand guidelines document "
                "spanning 84 pages.",
                font_size=14, font_color=RGBColor(0xCC, 0xCC, 0xCC))

    # --- Slide 5: Stripe Dashboard ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_bg(slide5, 0x23, 0x23, 0x23)
    add_textbox(slide5, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
                "Stripe — Dashboard Redesign", font_size=28, bold=True)
    add_textbox(slide5, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.5),
                "Redesigned the merchant analytics dashboard to improve data "
                "readability and reduce time-to-insight by 40%. Introduced a "
                "modular card system, customizable KPI widgets, and an AI-powered "
                "anomaly detection panel. Shipped to 2.3M merchants in Q3 2025.",
                font_size=14, font_color=RGBColor(0xCC, 0xCC, 0xCC))

    # --- Slide 6: Figma Plugin ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_bg(slide6, 0x1E, 0x1E, 0x1E)
    add_textbox(slide6, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
                "ColorSync — Figma Plugin", font_size=28, bold=True)
    add_textbox(slide6, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.5),
                "Designed and developed a Figma plugin that automatically "
                "syncs brand colors across team files. Used by 15,000+ designers "
                "with a 4.8-star rating. Features include palette generation from "
                "images, accessibility contrast checking, and theme switching.",
                font_size=14, font_color=RGBColor(0xCC, 0xCC, 0xCC))

    # --- Slide 7: Process ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_bg(slide7, 0x23, 0x23, 0x23)
    add_textbox(slide7, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
                "My Design Process", font_size=28, bold=True)
    steps = [
        ("1. Discover", "Stakeholder interviews, competitive audits, user research"),
        ("2. Define", "Problem framing, personas, journey mapping"),
        ("3. Design", "Wireframes, high-fidelity mockups, prototyping"),
        ("4. Deliver", "Developer handoff, QA, iteration based on metrics"),
    ]
    y_pos = 1.8
    for title, desc in steps:
        add_textbox(slide7, Inches(0.8), Inches(y_pos), Inches(10), Inches(0.5),
                    title, font_size=18, bold=True,
                    font_color=RGBColor(0x64, 0xB5, 0xF6))
        add_textbox(slide7, Inches(0.8), Inches(y_pos + 0.5), Inches(10), Inches(0.5),
                    desc, font_size=13, font_color=RGBColor(0xAA, 0xAA, 0xAA))
        y_pos += 1.2

    # --- Slide 8: Contact ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_bg(slide8, 0x1A, 0x1A, 0x2E)
    add_textbox(slide8, Inches(2.0), Inches(2.5), Inches(9), Inches(1.0),
                "Let's Work Together", font_size=36, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide8, Inches(2.0), Inches(4.0), Inches(9), Inches(0.8),
                "elena.vasquez@designstudio.co  |  +1 (415) 309-7742",
                font_size=16, font_color=RGBColor(0xBB, 0xBB, 0xBB),
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide8, Inches(2.0), Inches(5.0), Inches(9), Inches(0.6),
                "www.elenavasquez.design  |  @elenadesigns",
                font_size=14, font_color=RGBColor(0x64, 0xB5, 0xF6),
                alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
