"""
Reward Script: Image gallery slide with 6 placeholders in 3x2 grid
Task ID: impress_ps_026
Domain: libreoffice_impress
Scoring:
  Component 1: 6 rectangle placeholders exist on slide 3 (0.25)
  Component 2: Rectangles have white border ~1pt (0.20)
  Component 3: Rectangles arranged in 3x2 grid with consistent spacing (0.25)
  Component 4: 6 caption text boxes with 'Caption' text below rectangles (0.20)
  Component 5: Caption font is ~10pt and white color (0.10)
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_026'


def persist_app_state(domain):
    """Save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            import time
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def check_white_border(rect):
    """Check if a rectangle has a white border ~1pt. Returns 1 if yes, 0 if no."""
    try:
        line = rect.line
        if line.fill.type != 1:  # not SOLID
            return 0
        color = str(line.color.rgb).upper()
        width = line.width
        # 1pt = 12700 EMU; allow some tolerance
        if color == "FFFFFF" and width is not None and abs(width - 12700) <= 5000:
            return 1
    except Exception:
        pass
    return 0


def check_caption_font(cap):
    """Check if a caption text box has ~10pt white font. Returns 1 if yes, 0 if no."""
    for para in cap.text_frame.paragraphs:
        for run in para.runs:
            # 10pt = 127000 EMU; allow tolerance (8pt-12pt = 101600-152400)
            size_match = (run.font.size is not None and 101600 <= run.font.size <= 152400)
            color_match = False
            try:
                if run.font.color.type is not None and str(run.font.color.rgb).upper() == "FFFFFF":
                    color_match = (run.font.color.type is not None)
            except Exception:
                pass
            if size_match and color_match:
                return 1
    return 0


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[2]  # Slide 3 (0-indexed)

    # Collect rectangles (AUTO_SHAPE with solid fill) and caption text boxes on slide 3
    # Filter: only count AUTO_SHAPE with solid fill as image placeholders
    # (The title box on initial_env is AUTO_SHAPE but has BACKGROUND fill type)
    rectangles = []
    text_boxes = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # Only count shapes with solid fill as placeholder rectangles
            if shape.fill.type == 1:  # SOLID fill
                rectangles.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt.lower() == "caption":
                    text_boxes.append(shape)

    print(f"INFO: Found {len(rectangles)} rectangles and {len(text_boxes)} caption text boxes on slide 3")

    # Component 1: 6 rectangle placeholders exist on slide 3 (0.25 points)
    try:
        if len(rectangles) == 6:
            print(f"PASS: Component 1 — Found exactly 6 rectangle placeholders (0.25 pts)")
            total_score += 0.25
        elif len(rectangles) >= 4:
            partial = 0.25 * (len(rectangles) / 6.0)
            print(f"PARTIAL: Component 1 — Found {len(rectangles)}/6 rectangles ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — Expected 6 rectangles, found {len(rectangles)}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Rectangles have white border approximately 1pt (0.20 points)
    try:
        if len(rectangles) == 0:
            print(f"FAIL: Component 2 — No rectangles to check borders")
        else:
            border_ok_count = sum(check_white_border(rect) for rect in rectangles)

            if border_ok_count == 6:
                print(f"PASS: Component 2 — All 6 rectangles have white ~1pt border (0.20 pts)")
                total_score += 0.20
            elif border_ok_count >= 1:
                partial = 0.20 * (border_ok_count / 6.0)
                print(f"PARTIAL: Component 2 — {border_ok_count}/6 rectangles have correct border ({partial:.2f} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 2 — No rectangles have white ~1pt border")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Rectangles arranged in 3x2 grid with consistent spacing (0.25 points)
    try:
        if len(rectangles) < 6:
            print(f"FAIL: Component 3 — Need 6 rectangles for grid check, found {len(rectangles)}")
        else:
            # Sort rectangles by position: first by top (row), then by left (column)
            rects_sorted = sorted(rectangles, key=lambda s: (s.top, s.left))

            # Extract positions
            tops = [r.top for r in rects_sorted]
            lefts = [r.left for r in rects_sorted]

            # Check 2 distinct rows and 3 distinct columns
            unique_tops = sorted(set(tops))
            unique_lefts = sorted(set(lefts))

            grid_score = 0.0

            # Sub-check: 2 rows and 3 columns
            if len(unique_tops) == 2 and len(unique_lefts) == 3:
                print(f"  Grid structure: 3 columns x 2 rows detected")
                grid_score += 0.10

                # Sub-check: rows have 3 items each
                row1 = [r for r in rects_sorted if r.top == unique_tops[0]]
                row2 = [r for r in rects_sorted if r.top == unique_tops[1]]
                if len(row1) == 3 and len(row2) == 3:
                    grid_score += 0.05
                    print(f"  Row distribution: 3+3 correct")

                # Sub-check: consistent column spacing
                col_gaps = [unique_lefts[i+1] - unique_lefts[i] for i in range(len(unique_lefts)-1)]
                if len(col_gaps) >= 2:
                    gap_diff = abs(col_gaps[0] - col_gaps[1])
                    rel_diff = gap_diff / max(col_gaps) if max(col_gaps) > 0 else 0
                    if rel_diff <= 0.05:
                        grid_score += 0.05
                        print(f"  Column spacing consistent: gaps={col_gaps}")
                    else:
                        print(f"  Column spacing inconsistent: gaps={col_gaps}")

                # Sub-check: consistent rectangle sizes
                widths = [r.width for r in rectangles]
                heights = [r.height for r in rectangles]
                if max(widths) - min(widths) <= 50000 and max(heights) - min(heights) <= 50000:
                    grid_score += 0.05
                    print(f"  Rectangle sizes consistent: w~{widths[0]}, h~{heights[0]}")
                else:
                    print(f"  Rectangle sizes vary: w={min(widths)}-{max(widths)}, h={min(heights)}-{max(heights)}")
            else:
                print(f"  FAIL: Expected 2 rows and 3 cols, got {len(unique_tops)} rows and {len(unique_lefts)} cols")

            if grid_score > 0:
                print(f"PASS: Component 3 — Grid layout verified ({grid_score:.2f} pts)")
                total_score += grid_score
            else:
                print(f"FAIL: Component 3 — Grid layout not detected")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: 6 caption text boxes with 'Caption' text below rectangles (0.20 points)
    try:
        if len(text_boxes) == 6:
            # Verify captions are positioned below rectangles
            below_count = 0
            for cap in text_boxes:
                cap_top = cap.top
                cap_left = cap.left
                for rect in rectangles:
                    # Caption should be horizontally aligned with a rectangle and below it
                    left_match = abs(cap_left - rect.left) <= 50000
                    below = cap_top >= rect.top + rect.height - 50000
                    if left_match and below:
                        below_count += 1
                        break

            if below_count == 6:
                print(f"PASS: Component 4 — All 6 caption boxes positioned below rectangles (0.20 pts)")
                total_score += 0.20
            elif below_count >= 1:
                partial = 0.20 * (below_count / 6.0)
                print(f"PARTIAL: Component 4 — {below_count}/6 captions below rectangles ({partial:.2f} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 4 — No captions properly positioned below rectangles")
        elif len(text_boxes) >= 4:
            partial = 0.20 * (len(text_boxes) / 6.0)
            print(f"PARTIAL: Component 4 — Found {len(text_boxes)}/6 caption boxes ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — Expected 6 caption text boxes, found {len(text_boxes)}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Caption font is approximately 10pt and white color (0.10 points)
    try:
        if len(text_boxes) == 0:
            print(f"FAIL: Component 5 — No caption text boxes to check font")
        else:
            font_ok_count = sum(check_caption_font(cap) for cap in text_boxes)

            if font_ok_count == 6:
                print(f"PASS: Component 5 — All 6 captions have ~10pt white font (0.10 pts)")
                total_score += 0.10
            elif font_ok_count >= 1:
                partial = 0.10 * (font_ok_count / 6.0)
                print(f"PARTIAL: Component 5 — {font_ok_count}/6 captions have correct font ({partial:.2f} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 5 — No captions have ~10pt white font")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
