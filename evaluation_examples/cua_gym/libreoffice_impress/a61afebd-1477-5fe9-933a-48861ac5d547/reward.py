"""
Reward Script: Build comparative analysis presentation section (slides 4-7)
Task ID: impress_stu_073
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Slide 4 has a comparison table (7 rows x 3 cols) with correct headers
  Component 2 (0.25): Slide 5 has decision tree flowchart (multiple shapes including connectors and boxes)
  Component 3 (0.25): Slide 6 has two pie chart representations (2015 vs 2025 labels, oval shapes, legend entries)
  Component 4 (0.20): Slide 7 has summary content + styled quote box (Rounded Rectangle with quote text)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_073'


def count_shapes_by_type(slide):
    """Count shapes by their type on a slide."""
    counts = {}
    for shape in slide.shapes:
        stype = str(shape.shape_type)
        counts[stype] = counts.get(stype, 0) + 1
    return counts


def get_all_text(slide):
    """Get all text from a slide, concatenated."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return texts


def has_table(slide):
    """Check if slide has a table shape."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return True
    return False


def get_table(slide):
    """Get the first table on a slide."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return shape.table
    return None


def count_auto_shapes(slide):
    """Count auto shapes (rounded rectangles, ovals, etc.) on a slide."""
    count = 0
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            count += 1
    return count


def count_lines(slide):
    """Count line/connector shapes on a slide."""
    count = 0
    for shape in slide.shapes:
        # LINE type is 9
        if shape.shape_type == 9:  # MSO_SHAPE_TYPE.LINE
            count += 1
    return count


def has_rounded_rectangle(slide):
    """Check if slide has a Rounded Rectangle auto shape."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if 'Rounded Rectangle' in shape.name:
                return True
    return False


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Must have 10 slides
    if len(prs.slides) != 10:
        print(f"PRECONDITION FAIL: Expected 10 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)

    # =========================================================================
    # Component 1: Slide 4 — Comparison table (7 rows x 3 cols) (0.30 points)
    # Initial: Slide 4 has only 1 shape (title). Golden: has a TABLE shape.
    # =========================================================================
    try:
        slide4 = slides[3]  # 0-indexed
        table = get_table(slide4)
        if table is not None:
            nrows = len(table.rows)
            ncols = len(table.columns)
            print(f"  Slide 4 table found: {nrows}x{ncols}")

            # Check table dimensions: 7 rows (header + 6 criteria), 3 columns
            if nrows >= 7 and ncols == 3:
                # Check headers contain relevant text
                h0 = table.cell(0, 0).text.strip().lower()
                h1 = table.cell(0, 1).text.strip().lower()
                h2 = table.cell(0, 2).text.strip().lower()
                has_criteria = 'criteria' in h0 or 'criterion' in h0 or 'aspect' in h0
                has_qual = 'qualitative' in h1 or 'qualitative' in h2
                has_quant = 'quantitative' in h1 or 'quantitative' in h2

                if has_qual and has_quant:
                    # Full marks: correct table with both method columns
                    if nrows == 7 and has_criteria:
                        print(f"PASS: Component 1 — Slide 4 has 7x3 comparison table with correct headers (0.30 pts)")
                        total_score += 0.30
                    else:
                        # Partial: table present with qual/quant but dimensions or header slightly off
                        print(f"PARTIAL: Component 1 — Table has {nrows} rows, criteria header: {has_criteria} (0.20 pts)")
                        total_score += 0.20
                else:
                    print(f"PARTIAL: Component 1 — Table found but headers don't match: [{h0}] [{h1}] [{h2}] (0.10 pts)")
                    total_score += 0.10
            else:
                print(f"PARTIAL: Component 1 — Table found but wrong dimensions: {nrows}x{ncols} (0.10 pts)")
                total_score += 0.10
        else:
            # Check if slide 4 has more than just the title (some content added)
            num_shapes = len(slide4.shapes)
            if num_shapes > 1:
                print(f"FAIL: Component 1 — Slide 4 has {num_shapes} shapes but no table")
            else:
                print(f"FAIL: Component 1 — Slide 4 has only the title, no table added")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # =========================================================================
    # Component 2: Slide 5 — Decision tree flowchart (0.25 points)
    # Initial: 1 shape (title). Golden: 21 shapes (boxes, lines, text).
    # We check for: auto shapes (rounded rectangles) + line connectors + text content.
    # =========================================================================
    try:
        slide5 = slides[4]  # 0-indexed
        num_shapes = len(slide5.shapes)
        num_auto = count_auto_shapes(slide5)
        num_lines = count_lines(slide5)
        all_text = get_all_text(slide5)
        all_text_lower = ' '.join(all_text).lower()

        # A real flowchart needs: multiple auto shapes + connector lines
        has_flowchart_shapes = num_auto >= 5 and num_lines >= 3
        # Content should mention qualitative/quantitative/mixed methods
        has_qual_text = 'qualitative' in all_text_lower or 'qualitat' in all_text_lower
        has_quant_text = 'quantitative' in all_text_lower or 'quantitat' in all_text_lower
        has_method_content = has_qual_text and has_quant_text

        if has_flowchart_shapes and has_method_content:
            print(f"PASS: Component 2 — Slide 5 has flowchart ({num_auto} auto shapes, {num_lines} lines) with method content (0.25 pts)")
            total_score += 0.25
        elif num_shapes > 3 and has_method_content:
            # Some shapes added with right content but not full flowchart structure
            print(f"PARTIAL: Component 2 — Slide 5 has {num_shapes} shapes with method content but weak flowchart (0.15 pts)")
            total_score += 0.15
        elif num_shapes > 1:
            # Some content added but minimal
            print(f"PARTIAL: Component 2 — Slide 5 has {num_shapes} shapes but insufficient flowchart structure (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 2 — Slide 5 has only {num_shapes} shape(s), no flowchart")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # =========================================================================
    # Component 3: Slide 6 — Two pie chart representations for 2015 vs 2025 (0.25 points)
    # Initial: 1 shape (title). Golden: 20 shapes (ovals, rectangles for legend, text labels).
    # We check for: oval shapes (pie representations), year labels (2015, 2025), percentage data.
    # =========================================================================
    try:
        slide6 = slides[5]  # 0-indexed
        num_shapes = len(slide6.shapes)
        num_auto = count_auto_shapes(slide6)
        all_text = get_all_text(slide6)
        all_text_joined = ' '.join(all_text)

        has_2015 = '2015' in all_text_joined
        has_2025 = '2025' in all_text_joined
        has_both_years = has_2015 and has_2025

        # Check for percentage data in text
        has_percentages = '%' in all_text_joined

        # Check for oval shapes (pie chart representation)
        has_ovals = False
        oval_count = 0
        for shape in slide6.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and 'Oval' in shape.name:
                oval_count += 1
                has_ovals = True

        if has_both_years and has_ovals and num_shapes >= 10:
            if has_percentages and oval_count >= 2:
                print(f"PASS: Component 3 — Slide 6 has pie charts for both years ({oval_count} ovals, {num_shapes} shapes) with percentages (0.25 pts)")
                total_score += 0.25
            else:
                print(f"PARTIAL: Component 3 — Slide 6 has year labels and ovals but limited data (0.15 pts)")
                total_score += 0.15
        elif num_shapes > 3 and (has_2015 or has_2025):
            print(f"PARTIAL: Component 3 — Slide 6 has {num_shapes} shapes with some year data (0.10 pts)")
            total_score += 0.10
        elif num_shapes > 1:
            print(f"PARTIAL: Component 3 — Slide 6 has {num_shapes} shapes but missing pie chart structure (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 3 — Slide 6 has only {num_shapes} shape(s), no pie charts")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # =========================================================================
    # Component 4: Slide 7 — Summary + styled quote box (0.20 points)
    # Initial: 1 shape (title). Golden: 4 shapes (title, summary text, rounded rect quote box, footer text).
    # We check for: Rounded Rectangle shape (quote box), quote-related text content, multiple shapes.
    # =========================================================================
    try:
        slide7 = slides[6]  # 0-indexed
        num_shapes = len(slide7.shapes)
        all_text = get_all_text(slide7)
        all_text_joined = ' '.join(all_text).lower()

        has_quote_box = has_rounded_rectangle(slide7)
        # Check for quote content (mixed methods / Creswell or similar)
        has_quote_content = ('mixed method' in all_text_joined or 'qualitative' in all_text_joined
                             or 'quantitative' in all_text_joined)
        # Check for quotation marks or attribution
        has_attribution = ('"' in ' '.join(all_text) or '\u201c' in ' '.join(all_text)
                           or 'creswell' in all_text_joined or '- ' in ' '.join(all_text))

        if has_quote_box and has_quote_content and num_shapes >= 3:
            if has_attribution:
                print(f"PASS: Component 4 — Slide 7 has styled quote box with attribution ({num_shapes} shapes) (0.20 pts)")
                total_score += 0.20
            else:
                print(f"PARTIAL: Component 4 — Slide 7 has quote box and content but no clear attribution (0.15 pts)")
                total_score += 0.15
        elif num_shapes >= 3 and has_quote_content:
            print(f"PARTIAL: Component 4 — Slide 7 has content but no styled quote box (0.10 pts)")
            total_score += 0.10
        elif num_shapes > 1:
            print(f"PARTIAL: Component 4 — Slide 7 has {num_shapes} shapes but insufficient content (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 4 — Slide 7 has only {num_shapes} shape(s), no summary or quote")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist app state before verification (LibreOffice may have unsaved edits)
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
persist_app_state()
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
