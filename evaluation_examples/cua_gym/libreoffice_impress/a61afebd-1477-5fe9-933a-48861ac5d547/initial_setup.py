"""
Initial Setup: Build a comparative analysis presentation section (slides 4-7) with titles only
Task ID: impress_stu_073
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_073'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

TITLE_COLOR = RGBColor(0x2C, 0x3E, 0x50)
BG_COLOR = RGBColor(0xEC, 0xF0, 0xF1)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, size=Pt(32)):
    """Add a title textbox at the top of a blank slide."""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = size
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR
    return txBox


def add_body_text(slide, text, left=0.5, top=1.5, width=9.0, height=4.5, size=Pt(16)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = size
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    return txBox


def add_bullet_points(slide, items, left=0.5, top=1.5, width=9.0, height=5.0, size=Pt(14)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.space_after = Pt(6)
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = size
        run.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    s1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(s1, BG_COLOR)
    add_title(s1, "Research Methods in Social Sciences", size=Pt(36))
    add_body_text(s1, "A Comprehensive Overview for Graduate Students\nDr. Elena Vasquez | Spring 2025",
                  top=2.5, size=Pt(20))

    # --- Slide 2: Course Overview ---
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s2, BG_COLOR)
    add_title(s2, "Course Overview")
    add_bullet_points(s2, [
        "Introduction to research paradigms and philosophies",
        "Understanding qualitative and quantitative approaches",
        "Mixed methods research design",
        "Data collection techniques across disciplines",
        "Ethical considerations in modern research",
        "Publishing and peer review processes"
    ])

    # --- Slide 3: Research Paradigms ---
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s3, BG_COLOR)
    add_title(s3, "Research Paradigms")
    add_bullet_points(s3, [
        "Positivism: Objective reality, measurable phenomena",
        "Constructivism: Reality is socially constructed",
        "Pragmatism: Focus on practical outcomes",
        "Critical Theory: Power dynamics and social justice",
        "Post-positivism: Modified objectivity with acknowledgment of bias"
    ])

    # --- Slide 4: Title Only (Comparative Analysis) ---
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s4, BG_COLOR)
    add_title(s4, "Comparative Analysis: Qualitative vs Quantitative")

    # --- Slide 5: Title Only (Decision Tree) ---
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s5, BG_COLOR)
    add_title(s5, "When to Use Each Method")

    # --- Slide 6: Title Only (Pie Charts) ---
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s6, BG_COLOR)
    add_title(s6, "Distribution of Methods in Published Research")

    # --- Slide 7: Title Only (Summary Quote) ---
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s7, BG_COLOR)
    add_title(s7, "Summary & Key Quote")

    # --- Slide 8: Data Collection Methods ---
    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s8, BG_COLOR)
    add_title(s8, "Data Collection Methods")
    add_bullet_points(s8, [
        "Surveys and questionnaires (structured, semi-structured)",
        "In-depth interviews and focus groups",
        "Observation (participant and non-participant)",
        "Document analysis and archival research",
        "Experiments and quasi-experiments",
        "Case studies and ethnographic methods"
    ])

    # --- Slide 9: Ethical Considerations ---
    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s9, BG_COLOR)
    add_title(s9, "Ethical Considerations in Research")
    add_bullet_points(s9, [
        "Informed consent and voluntary participation",
        "Confidentiality and data protection (GDPR compliance)",
        "Avoiding harm to participants",
        "Institutional Review Board (IRB) approval",
        "Responsible data management and storage",
        "Transparency in reporting findings"
    ])

    # --- Slide 10: References & Further Reading ---
    s10 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s10, BG_COLOR)
    add_title(s10, "References & Further Reading")
    add_bullet_points(s10, [
        "Creswell, J.W. (2023). Research Design: Qualitative, Quantitative, and Mixed Methods",
        "Bryman, A. (2022). Social Research Methods, 6th Edition",
        "Tashakkori, A. & Teddlie, C. (2021). Foundations of Mixed Methods Research",
        "Denzin, N.K. & Lincoln, Y.S. (2024). The SAGE Handbook of Qualitative Research",
        "Field, A. (2024). Discovering Statistics Using IBM SPSS Statistics, 6th Edition"
    ], size=Pt(12))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


def launch_gui(command: str, delay_sec: float = 1.0):
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


create_initial()
