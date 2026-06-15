"""
Initial Setup: Configure custom slide shows in a 15-slide master deck
Task ID: impress_gf3_021
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_021'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    """Add a title slide (layout 0)."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points):
    """Add a title + content slide (layout 1)."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.paragraphs[0].text = point
        else:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide (Intro section) ---
    add_title_slide(
        prs,
        "Nextera Analytics Platform",
        "Q2 2025 Strategic Review\nPrepared by Sarah Chen, VP of Engineering"
    )

    # --- Slide 2: Company Overview (Intro section) ---
    add_content_slide(prs, "Company Overview", [
        "Founded in 2019 by former Google and Amazon engineers",
        "Headquarters in San Francisco, with offices in London and Singapore",
        "450+ employees across engineering, product, and sales",
        "Series C funded: $128M raised to date",
        "Serving 2,300+ enterprise clients in 45 countries",
    ])

    # --- Slide 3: Mission & Vision (Intro section) ---
    add_content_slide(prs, "Mission & Vision", [
        "Mission: Democratize data analytics for mid-market enterprises",
        "Vision: Become the industry-standard platform for real-time decision intelligence",
        "Core Values: Transparency, Innovation, Customer Obsession",
        "2025 Theme: 'Scale with Purpose'",
    ])

    # --- Slide 4: System Architecture (Technical section) ---
    add_content_slide(prs, "System Architecture", [
        "Microservices architecture with 47 independent services",
        "Event-driven communication via Apache Kafka (3.2M events/sec peak)",
        "Kubernetes orchestration across 3 AWS regions",
        "Service mesh with Istio for traffic management and observability",
        "GitOps deployment pipeline with ArgoCD",
    ])

    # --- Slide 5: Backend Infrastructure (Technical section) ---
    add_content_slide(prs, "Backend Infrastructure", [
        "Primary data store: PostgreSQL 16 with Citus extension (horizontal sharding)",
        "Real-time cache: Redis Cluster with 256GB capacity",
        "Search engine: Elasticsearch 8.x with 12-node cluster",
        "Message queue: RabbitMQ for task orchestration",
        "Object storage: MinIO (S3-compatible) for raw data lake",
    ])

    # --- Slide 6: Frontend Platform (Technical section) ---
    add_content_slide(prs, "Frontend Platform", [
        "React 18 with TypeScript - 380K lines of code",
        "Micro-frontend architecture using Module Federation",
        "Design system: 'Nextera UI' with 120+ reusable components",
        "Performance: LCP < 1.2s, FID < 50ms across all dashboards",
        "Accessibility: WCAG 2.1 AA compliant",
    ])

    # --- Slide 7: Database Layer (Technical section) ---
    add_content_slide(prs, "Database Layer & Data Pipeline", [
        "ETL pipeline processing 2.8TB daily using Apache Spark",
        "Data warehouse: Snowflake with 45 curated schemas",
        "Real-time streaming: Flink jobs for sub-second aggregations",
        "Data quality framework: Great Expectations with 1,200+ validation rules",
        "Backup: Continuous replication with 15-minute RPO",
    ])

    # --- Slide 8: API Design (Technical section) ---
    add_content_slide(prs, "API Design & Integration", [
        "RESTful API v3 with OpenAPI 3.1 specification",
        "GraphQL gateway for flexible client queries",
        "Rate limiting: 10K requests/min per tenant",
        "OAuth 2.0 + OIDC for authentication, RBAC for authorization",
        "SDK support: Python, JavaScript, Go, Java",
    ])

    # --- Slide 9: Testing Strategy (Technical section) ---
    add_content_slide(prs, "Testing Strategy", [
        "Unit test coverage: 89% across all services",
        "Integration tests: 2,400+ end-to-end scenarios",
        "Performance testing: weekly load tests simulating 50K concurrent users",
        "Chaos engineering: monthly GameDay exercises",
        "Security: quarterly penetration testing by NCC Group",
    ])

    # --- Slide 10: Deployment & DevOps (Technical section) ---
    add_content_slide(prs, "Deployment & DevOps", [
        "CI/CD: GitHub Actions with 12-minute average pipeline time",
        "Deploy frequency: 25+ production deployments per week",
        "Feature flags: LaunchDarkly managing 340 active flags",
        "Canary deployments with automated rollback on error spike",
        "Infrastructure as Code: Terraform managing 1,800+ resources",
    ])

    # --- Slide 11: Performance Metrics (Technical section) ---
    add_content_slide(prs, "Performance & Reliability", [
        "Uptime: 99.97% over last 12 months (13 minutes total downtime)",
        "P99 API latency: 180ms for read operations, 320ms for writes",
        "Auto-scaling: handles 3x traffic spikes within 90 seconds",
        "MTTR: average 8 minutes for P1 incidents",
        "Error budget: consuming 23% of quarterly allocation",
    ])

    # --- Slide 12: Market Analysis (Business section) ---
    add_content_slide(prs, "Market Analysis", [
        "TAM: $47.2B (enterprise analytics market, Gartner 2025)",
        "Current market share: 3.8% with trajectory to 6.2% by 2027",
        "Key competitors: Tableau, Looker, Power BI, Domo",
        "Competitive advantage: Real-time processing + self-service ML",
        "Customer NPS: 72 (industry average: 41)",
    ])

    # --- Slide 13: Revenue Model (Business section) ---
    add_content_slide(prs, "Revenue Model & Growth", [
        "ARR: $86.4M (42% YoY growth)",
        "Net Revenue Retention: 135%",
        "Average deal size: $37,500/year (up from $28,000 in 2024)",
        "Enterprise tier (>$100K/yr): 18% of customers, 62% of revenue",
        "Gross margin: 78%",
    ])

    # --- Slide 14: Roadmap & Timeline (Business section) ---
    add_content_slide(prs, "2025 Product Roadmap", [
        "Q2: Launch AI-powered anomaly detection (Project Sentinel)",
        "Q3: Multi-cloud support (Azure, GCP alongside AWS)",
        "Q3: Natural language query interface (GPT-4 integration)",
        "Q4: Self-service ML model deployment for business users",
        "Q4: SOC 2 Type II and ISO 27001 certification completion",
    ])

    # --- Slide 15: Q&A (Business section) ---
    add_title_slide(
        prs,
        "Questions & Discussion",
        "Thank you for your time\nContact: sarah.chen@nextera-analytics.com"
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
