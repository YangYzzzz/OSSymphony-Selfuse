"""
FINAL REWARD SCRIPT - SUCCESS
Task: Align Table 1 to the left and set spacing to text = 0.25 cm.
Generated: 2025-10-17 09:34:42
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation


def verify_task(file_path: str) -> float:
    """Verify that:
    1. Table 1 is horizontally aligned to the extreme left of the slide
    2. Spacing to text (cell margins) is set to exactly 0.25 cm (≈ 90 000 EMU) on all sides

    Returns a progressive score between 0.0 and 1.0
    """
    print(f"Verifying presentation: {file_path}")

    # Scoring weights (total 1.0)
    ALIGN_WEIGHT = 0.5      # Left-alignment contributes 0.5
    SPACING_WEIGHT = 0.5    # Correct margins contribute 0.5 (0.125 each side)

    # ------------------------------------------------------------------
    # 1. File existence & loading (prerequisite – no points awarded)
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File not found")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error loading PPTX: {e}")
        return 0.0

    # ------------------------------------------------------------------
    # 2. Locate the first table (Table 1) in the presentation
    # ------------------------------------------------------------------
    table_shape = None
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table_shape = shape
                break
        if table_shape:
            break

    if table_shape is None:
        print("✗ No table found in the presentation")
        return 0.0
    print("✓ Found a table shape (assumed Table 1)")

    total_score = 0.0  # progressive scoring starts here

    # ------------------------------------------------------------------
    # 3. Verify left-alignment (x-position == 0 with small tolerance)
    # ------------------------------------------------------------------
    LEFT_TOLERANCE = 20_000  # ≈ 0.02 cm in EMU
    if table_shape.left <= LEFT_TOLERANCE:
        print(f"✓ Table is aligned left (left = {table_shape.left} EMU)")
        total_score += ALIGN_WEIGHT
    else:
        print(f"✗ Table not fully left-aligned (left = {table_shape.left} EMU)")

    # ------------------------------------------------------------------
    # 4. Verify spacing to text (cell margins) == 0.25 cm on ALL sides
    # ------------------------------------------------------------------
    TARGET_MARGIN = 90_000   # 0.25 cm in EMU
    MARGIN_TOL   = 2_000    # ±0.002 cm tolerance

    tbl = table_shape.table
    sample_cell = tbl.cell(0, 0)  # any cell works – margins are uniform per table

    margins = {
        "left":   sample_cell.margin_left,
        "right":  sample_cell.margin_right,
        "top":    sample_cell.margin_top,
        "bottom": sample_cell.margin_bottom,
    }

    correct_sides = 0
    for side, value in margins.items():
        if abs(value - TARGET_MARGIN) <= MARGIN_TOL:
            print(f"✓ Margin {side} ≈ 0.25 cm ({value} EMU)")
            correct_sides += 1
        else:
            print(f"✗ Margin {side} incorrect ({value} EMU, expected ≈ {TARGET_MARGIN})")

    # Award partial credit for each correctly-set side
    if correct_sides:
        total_score += SPACING_WEIGHT * (correct_sides / 4)

    # ------------------------------------------------------------------
    # 5. Final score (capped at 1.0)
    # ------------------------------------------------------------------
    final_score = min(total_score, 1.0)
    print(f"Final computed score: {final_score}")
    return final_score


if __name__ == "__main__":
    # Path to the presentation inside the VM workspace
    FILE_PATH = "/home/user/align_table_1_to_the_left_and_set_spacing_to_text_025_cm.pptx"

    reward = verify_task(FILE_PATH)
    print(f"REWARD: {reward}")
