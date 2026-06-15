"""
Reward Script: LibreOffice Impress macro to generate a summary slide
Task ID: impress_gf5_021
Domain: libreoffice_impress
Scoring:
  Component 1 (0.2): Presentation has 13 slides
  Component 2 (0.2): Last slide title is "Summary"
  Component 3 (0.3): Last slide content lists titles of slides 1-12 as bullet points
  Component 4 (0.3): Macro "AddSummarySlide" exists in ODP file
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_021'

# Expected titles of the original 12 slides
EXPECTED_TITLES = [
    'Executive Summary',
    'Q1 Revenue Analysis',
    'Market Segmentation Overview',
    'Customer Acquisition Metrics',
    'Product Roadmap 2025',
    'Competitive Landscape',
    'Operational Efficiency Report',
    'Team Performance Highlights',
    'Risk Assessment and Mitigation',
    'Financial Projections',
    'Strategic Partnerships',
    'Next Steps and Action Items',
]


def verify_task():
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    pptx_path = f'{WORKDIR}/results_presentation.pptx'
    odp_path = f'{WORKDIR}/results_presentation.odp'

    # Load presentation
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load PPTX file {pptx_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Presentation has 13 slides (0.2 points)
    # Initial env has 12 slides; golden should have 13 (12 + summary)
    try:
        if num_slides == 13:
            print(f"PASS: Component 1 -- Slide count is 13 (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 1 -- Expected 13 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Last slide title is "Summary" (0.2 points)
    # This should fail on initial env (no Summary slide)
    try:
        if num_slides >= 13:
            last_slide = prs.slides[num_slides - 1]
            last_title = ''
            for shape in last_slide.shapes:
                if shape.has_text_frame and shape.text_frame.paragraphs:
                    last_title = shape.text_frame.paragraphs[0].text.strip()
                    break
            if last_title.lower() == 'summary':
                print(f"PASS: Component 2 -- Last slide title is 'Summary' (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 2 -- Last slide title is '{last_title}', expected 'Summary'")
        else:
            print(f"FAIL: Component 2 -- Not enough slides to check last slide title")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Last slide content lists titles of slides 1-12 (0.3 points)
    # The content placeholder should have bullet points with all 12 original titles
    try:
        if num_slides >= 13:
            last_slide = prs.slides[num_slides - 1]
            # Find the content placeholder (second shape, typically)
            content_texts = []
            for idx, shape in enumerate(last_slide.shapes):
                if idx == 0:
                    continue  # skip title shape
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            content_texts.append(text)

            # Check how many of the 12 original titles appear in the bullet list
            matched = 0
            for expected_title in EXPECTED_TITLES:
                if any(expected_title.lower() in ct.lower() for ct in content_texts):
                    matched += 1

            if matched == 12:
                print(f"PASS: Component 3 -- All 12 slide titles found in summary content (0.3 pts)")
                total_score += 0.3
            elif matched >= 8:
                partial = round(0.3 * (matched / 12), 2)
                print(f"PARTIAL: Component 3 -- {matched}/12 titles found ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 3 -- Only {matched}/12 titles found in summary content")
        else:
            print(f"FAIL: Component 3 -- Not enough slides to check content")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Macro "AddSummarySlide" exists in ODP file (0.3 points)
    # Check the ODP ZIP for Basic/Standard/AddSummarySlide.xml
    try:
        if os.path.exists(odp_path):
            with zipfile.ZipFile(odp_path, 'r') as zf:
                names = zf.namelist()
                macro_found = any('AddSummarySlide' in n for n in names)
                if macro_found:
                    # Also verify it contains the macro code
                    macro_file = [n for n in names if 'AddSummarySlide' in n][0]
                    content = zf.read(macro_file).decode('utf-8', errors='replace')
                    # Check for key elements: Sub AddSummarySlide, Summary, getByIndex
                    has_sub = 'AddSummarySlide' in content
                    has_summary_text = 'Summary' in content
                    has_slide_access = 'getByIndex' in content or 'DrawPages' in content

                    if has_sub and has_summary_text and has_slide_access:
                        print(f"PASS: Component 4 -- Macro 'AddSummarySlide' found with correct structure (0.3 pts)")
                        total_score += 0.3
                    elif has_sub:
                        print(f"PARTIAL: Component 4 -- Macro exists but may be incomplete (0.15 pts)")
                        total_score += 0.15
                    else:
                        print(f"FAIL: Component 4 -- Macro file found but content invalid")
                else:
                    print(f"FAIL: Component 4 -- No 'AddSummarySlide' macro found in ODP")
        else:
            print(f"FAIL: Component 4 -- ODP file not found at {odp_path}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
pptx_path = f'{WORKDIR}/results_presentation.pptx'
if not os.path.exists(pptx_path):
    print(f"File not found: {pptx_path}")
    print("REWARD: 0.0")
else:
    verify_task()
