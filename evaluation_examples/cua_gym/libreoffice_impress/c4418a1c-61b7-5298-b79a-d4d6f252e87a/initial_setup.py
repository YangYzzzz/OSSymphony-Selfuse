"""
Initial Setup: Create a 12-slide presentation for the macro task
Task ID: impress_gf5_021
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
OUTPUT = f'{WORKDIR}/results_presentation.pptx'

SLIDE_TITLES = [
    "Executive Summary",
    "Q1 Revenue Analysis",
    "Market Segmentation Overview",
    "Customer Acquisition Metrics",
    "Product Roadmap 2025",
    "Competitive Landscape",
    "Operational Efficiency Report",
    "Team Performance Highlights",
    "Risk Assessment and Mitigation",
    "Financial Projections",
    "Strategic Partnerships",
    "Next Steps and Action Items",
]

SLIDE_CONTENTS = [
    "This presentation covers the key business results and strategic initiatives for the current fiscal year. All departments contributed data and analysis.",
    "Total revenue reached $4.2M in Q1, representing a 15% increase year-over-year. The APAC region showed the strongest growth at 23%.",
    "Our customer base is divided into Enterprise (45%), Mid-Market (30%), and SMB (25%) segments. Enterprise accounts drove 60% of total revenue.",
    "New customer acquisition cost decreased to $185 per customer. Monthly active users grew by 12,000, with a retention rate of 89%.",
    "Key milestones include the v3.0 platform release in June, mobile app redesign in August, and API gateway launch in October.",
    "Three main competitors identified: AlphaTech (35% market share), BetaSoft (20%), and GammaIO (15%). Our differentiation lies in integration capabilities.",
    "Automated deployment pipeline reduced release cycles from 2 weeks to 3 days. Infrastructure costs decreased by 18% through cloud optimization.",
    "Engineering team completed 94% of sprint commitments. Sales exceeded quarterly targets by 8%. Customer support resolved 97% of tickets within SLA.",
    "Top risks: supply chain disruption (medium), regulatory changes in EU markets (high), and talent retention in engineering (medium).",
    "Projected annual revenue of $18.5M with an EBITDA margin of 22%. Cash runway extends through Q2 2027 at current burn rate.",
    "Signed MOU with DataFlow Inc. for data integration. Exploring partnership with CloudNine for enterprise distribution in Latin America.",
    "Complete board presentation by April 15. Finalize hiring plan for Q2. Launch customer advisory board pilot program.",
]


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    for i, (title, content) in enumerate(zip(SLIDE_TITLES, SLIDE_CONTENTS)):
        # Use layout 1 = Title, Content
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # Set title
        slide.shapes.title.text = title
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            run.font.size = Pt(28)
            run.font.bold = True

        # Set content in content placeholder (index 1)
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = content
        for run in tf.paragraphs[0].runs:
            run.font.size = Pt(16)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Number of slides: {len(prs.slides)}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
