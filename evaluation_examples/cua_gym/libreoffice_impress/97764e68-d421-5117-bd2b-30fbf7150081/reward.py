"""
Reward Script: Split-screen testimonial layout on slide 6
Task ID: impress_sales_081
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): Left-side light blue background rectangle exists
  Component 2 (0.10): Oval customer photo placeholder exists
  Component 3 (0.15): Quote text in italic ~16pt exists
  Component 4 (0.10): Attribution text in bold ~12pt exists
  Component 5 (0.25): Three metric callout shapes with correct text and fill colors
  Component 6 (0.20): Fly In animations from opposite directions
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_081'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.oxml.ns import qn
        from pptx.util import Pt, Emu
        from pptx.dml.color import RGBColor
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[5]  # slide 6 (0-indexed)

    # Gather all shapes on slide 6 (excluding the pre-existing title/textbox)
    # We need shapes beyond the initial 2 (Title placeholder + "Customer Impact" textbox)
    shapes = list(slide.shapes)

    # Helper: get shape fill color as hex string
    def get_fill_color(shape):
        try:
            fill = shape.fill
            if fill.type is not None and fill.type == 1:  # SOLID
                return str(fill.fore_color.rgb).upper()
        except Exception:
            pass
        return None

    # Helper: get preset geometry type
    def get_preset_geom(shape):
        try:
            el = shape._element
            for child in el:
                if 'spPr' in child.tag:
                    for gc in child:
                        if 'prstGeom' in gc.tag:
                            return gc.get('prst')
        except Exception:
            pass
        return None

    # Helper: get all text from shape
    def get_shape_text(shape):
        if shape.has_text_frame:
            return shape.text_frame.text.strip()
        return ""

    # Helper: get font properties from first non-empty run
    def get_first_run_font(shape):
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        return run.font
        return None

    # =========================================================================
    # Component 1: Light blue (#E8F4FD) background rectangle on left side (0.20)
    # =========================================================================
    try:
        found_blue_rect = False
        mid_x = prs.slide_width // 2  # midpoint of slide

        for shape in shapes:
            color = get_fill_color(shape)
            geom = get_preset_geom(shape)
            # Must be a rectangle with E8F4FD fill, positioned in left half
            if color == 'E8F4FD' and geom == 'rect' and shape.left < mid_x:
                found_blue_rect = True
                break

        if found_blue_rect:
            print(f"PASS: Component 1 — Light blue (#E8F4FD) rectangle found on left side (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 — No light blue (#E8F4FD) rectangle found on left side of slide 6")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # =========================================================================
    # Component 2: Oval customer photo placeholder (0.10)
    # =========================================================================
    try:
        found_oval = False
        for shape in shapes:
            geom = get_preset_geom(shape)
            if geom == 'ellipse' and shape.left < mid_x:
                found_oval = True
                break

        if found_oval:
            print(f"PASS: Component 2 — Oval (ellipse) photo placeholder found on left side (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 — No oval shape found on left side of slide 6")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # =========================================================================
    # Component 3: Quote text in italic ~16pt (0.15)
    # =========================================================================
    try:
        found_quote = False
        for shape in shapes:
            text = get_shape_text(shape)
            # Quote should be a substantial text block with quotation marks or meaningful content
            if len(text) > 30 and shape.left < mid_x:
                font = get_first_run_font(shape)
                if font is not None:
                    is_italic = font.italic is True
                    # 16pt = 203200 EMU; allow tolerance of +/- 4pt
                    size_ok = font.size is not None and abs(font.size - 203200) < 60000
                    if is_italic and size_ok:
                        found_quote = True
                        print(f"  Detail: quote text found, italic={font.italic}, size={font.size}")
                        break

        if found_quote:
            print(f"PASS: Component 3 — Quote text in italic ~16pt found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 3 — No italic ~16pt quote text found on left side of slide 6")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # =========================================================================
    # Component 4: Attribution text in bold ~12pt (0.10)
    # =========================================================================
    try:
        found_attribution = False
        for shape in shapes:
            text = get_shape_text(shape)
            # Attribution is typically shorter, contains a dash or name
            if shape.left < mid_x and 5 < len(text) < 200:
                font = get_first_run_font(shape)
                if font is not None:
                    is_bold = font.bold is True
                    # 12pt = 152400 EMU; allow tolerance of +/- 4pt
                    size_ok = font.size is not None and abs(font.size - 152400) < 60000
                    # Distinguish from quote: must be bold and not italic
                    is_not_italic = font.italic is not True
                    if is_bold and size_ok and is_not_italic:
                        found_attribution = True
                        print(f"  Detail: attribution text='{text[:50]}', bold={font.bold}, size={font.size}")
                        break

        if found_attribution:
            print(f"PASS: Component 4 — Attribution text in bold ~12pt found (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 4 — No bold ~12pt attribution text found on left side of slide 6")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # =========================================================================
    # Component 5: Three metric callout shapes with correct text/colors (0.25)
    # Points: ~0.083 per metric shape found correctly
    # =========================================================================
    try:
        metric_score = 0.0
        expected_metrics = [
            {'text_contains': 'Revenue', 'text_contains2': '+40%', 'color_family': 'green'},
            {'text_contains': 'Costs', 'text_contains2': '-25%', 'color_family': 'blue'},
            {'text_contains': 'NPS', 'text_contains2': '+30', 'color_family': 'orange'},
        ]

        def is_color_family(hex_color, family):
            """Check if a hex color belongs to a general color family."""
            if hex_color is None:
                return False
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            if family == 'green':
                return g > r and g > b  # green channel dominant
            elif family == 'blue':
                return b > r and b > g  # blue channel dominant
            elif family == 'orange':
                return r > g and r > b and g > b  # red dominant, some green
            return False

        metrics_found = 0
        for metric in expected_metrics:
            for shape in shapes:
                text = get_shape_text(shape).lower()
                if metric['text_contains'].lower() in text and metric['text_contains2'].lower() in text:
                    # Check that shape is on right half
                    if shape.left >= mid_x:
                        fill_color = get_fill_color(shape)
                        if fill_color and is_color_family(fill_color, metric['color_family']):
                            metrics_found += 1
                            print(f"  Detail: Metric '{metric['text_contains']}' found with {metric['color_family']} fill ({fill_color})")
                            break

        metric_score = (metrics_found / 3.0) * 0.25
        if metrics_found == 3:
            print(f"PASS: Component 5 — All 3 metric callout shapes found with correct colors (0.25 pts)")
        elif metrics_found > 0:
            print(f"PARTIAL: Component 5 — {metrics_found}/3 metric shapes found ({metric_score:.3f} pts)")
        else:
            print(f"FAIL: Component 5 — No metric callout shapes found on right side of slide 6")
        total_score += metric_score
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # =========================================================================
    # Component 6: Fly In animations from opposite directions (0.20)
    # Left side shapes should fly in from left (subtype 8)
    # Right side shapes should fly in from a different direction (subtype != 8)
    # =========================================================================
    try:
        animation_score = 0.0
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

        # Build shape_id -> position mapping
        shape_id_to_side = {}
        for shape in shapes:
            sid = str(shape.shape_id)
            if shape.left < mid_x:
                shape_id_to_side[sid] = 'left'
            else:
                shape_id_to_side[sid] = 'right'

        with zipfile.ZipFile(file_path, 'r') as zf:
            try:
                with zf.open('ppt/slides/slide6.xml') as f:
                    root = ET.parse(f).getroot()

                    # Find all Fly In animations (presetID=2, presetClass=entr)
                    left_anims = []
                    right_anims = []

                    for cTn in root.iter('{http://schemas.openxmlformats.org/presentationml/2006/main}cTn'):
                        pid = cTn.get('presetID')
                        pclass = cTn.get('presetClass')
                        psub = cTn.get('presetSubtype')
                        if pid == '2' and pclass == 'entr':
                            # Find target spid
                            for spTgt in cTn.iter('{http://schemas.openxmlformats.org/presentationml/2006/main}spTgt'):
                                spid = spTgt.get('spid')
                                side = shape_id_to_side.get(spid, 'unknown')
                                if side == 'left':
                                    left_anims.append(psub)
                                elif side == 'right':
                                    right_anims.append(psub)

                    print(f"  Detail: Left-side Fly In subtypes: {left_anims}")
                    print(f"  Detail: Right-side Fly In subtypes: {right_anims}")

                    # Check: at least one left-side shape has Fly In animation
                    has_left_anim = len(left_anims) > 0
                    # Check: at least one right-side shape has Fly In animation
                    has_right_anim = len(right_anims) > 0
                    # Check: directions are different (opposite sides)
                    directions_differ = False
                    if has_left_anim and has_right_anim:
                        left_subtypes = set(left_anims)
                        right_subtypes = set(right_anims)
                        directions_differ = not left_subtypes.intersection(right_subtypes)

                    if has_left_anim and has_right_anim and directions_differ:
                        animation_score = 0.20
                        print(f"PASS: Component 6 — Fly In animations from opposite directions (0.20 pts)")
                    elif has_left_anim and has_right_anim:
                        animation_score = 0.10
                        print(f"PARTIAL: Component 6 — Fly In animations present but same direction (0.10 pts)")
                    elif has_left_anim or has_right_anim:
                        animation_score = 0.05
                        print(f"PARTIAL: Component 6 — Fly In animation only on one side (0.05 pts)")
                    else:
                        print(f"FAIL: Component 6 — No Fly In animations found on slide 6")

            except KeyError:
                print(f"FAIL: Component 6 — Could not open slide6.xml")

        total_score += animation_score
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
