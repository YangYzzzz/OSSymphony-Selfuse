"""
Initial Setup: Create a 9-slide sales presentation with slide 6 titled 'Customer Impact' but empty content.
Task ID: impress_sales_081
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_081'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    """Add a title slide (layout 0)."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_items):
    """Add a title + content slide (layout 1) with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, item in enumerate(bullet_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(18)
    return slide


def add_title_only_slide(prs, title_text):
    """Add a title-only slide (layout 5=blank, manually add title)."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    add_title_slide(
        prs,
        "Q3 2025 Sales Performance Review",
        "Prepared by the Revenue Operations Team | September 2025"
    )

    # --- Slide 2: Executive Summary ---
    add_content_slide(prs, "Executive Summary", [
        "Overall revenue grew 18% YoY to $12.4M in Q3",
        "Enterprise segment outperformed targets by 22%",
        "Customer acquisition costs decreased by 15%",
        "Net Promoter Score improved from 42 to 56",
        "Three new strategic partnerships signed",
    ])

    # --- Slide 3: Revenue Overview ---
    add_content_slide(prs, "Revenue Overview", [
        "Total Q3 Revenue: $12,432,500 (+18% YoY)",
        "Recurring Revenue: $8,215,000 (66% of total)",
        "New Business: $2,847,300 (+31% vs Q2)",
        "Expansion Revenue: $1,370,200 (+12% YoY)",
        "Average Deal Size: $47,500 (up from $39,200)",
    ])

    # --- Slide 4: Regional Performance ---
    add_content_slide(prs, "Regional Performance", [
        "North America: $6.2M (50% of revenue, +15% YoY)",
        "Europe: $3.1M (25%, +24% YoY) - strongest growth",
        "Asia-Pacific: $2.0M (16%, +19% YoY)",
        "Latin America: $1.1M (9%, +8% YoY)",
        "Key win: Deutsche Telekom ($380K annual contract)",
    ])

    # --- Slide 5: Product Line Analysis ---
    add_content_slide(prs, "Product Line Analysis", [
        "Platform Pro: $5.8M (+22%) - enterprise tier driving growth",
        "Platform Essentials: $3.9M (+11%) - stable mid-market",
        "API Suite: $1.8M (+45%) - fastest growing segment",
        "Professional Services: $0.9M (-5%) - transitioning to partner-led",
        "New: Platform Pro AI Module beta launched with 12 pilot customers",
    ])

    # --- Slide 6: Customer Impact (EMPTY content - task target) ---
    add_title_only_slide(prs, "Customer Impact")
    # NOTE: No additional shapes, no testimonial layout, no animations.
    # This is the slide the agent must populate.

    # --- Slide 7: Strategic Initiatives ---
    add_content_slide(prs, "Strategic Initiatives", [
        "Launch AI-powered analytics dashboard by Q4",
        "Expand partner ecosystem: 15 new certified partners",
        "Implement usage-based pricing for API Suite",
        "Open Tokyo office to accelerate APAC growth",
        "Invest $2M in customer success automation",
    ])

    # --- Slide 8: Team Highlights ---
    add_content_slide(prs, "Team Highlights", [
        "Sales team expanded to 45 reps (+8 this quarter)",
        "Top performer: Rachel Torres - $1.2M in closed deals",
        "Average ramp time reduced from 6 to 4 months",
        "Sales enablement NPS: 78 (up from 65)",
        "12 team members completed advanced certification",
    ])

    # --- Slide 9: Next Steps & Q4 Outlook ---
    add_content_slide(prs, "Next Steps & Q4 Outlook", [
        "Target: $14.5M revenue in Q4 (17% growth)",
        "Priority: Close 5 enterprise deals in pipeline ($2.1M total)",
        "Launch customer referral program by October 15",
        "Complete CRM migration to unified platform",
        "Schedule quarterly business reviews with top 20 accounts",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
