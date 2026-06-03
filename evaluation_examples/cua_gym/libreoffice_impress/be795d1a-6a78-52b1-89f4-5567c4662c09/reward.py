"""
Reward Script: Add financial scorecard table to slide 4
Task ID: impress_exec_056
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.25): Table exists on slide 4 with correct dimensions (5x6)
  - Component 2 (0.25): Header row text matches expected values
  - Component 3 (0.25): All 20 data cell values match expected values
  - Component 4 (0.15): Header row fill #003366 with white text
  - Component 5 (0.10): Data cells use 14pt Calibri font
"""

import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Pt

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_056'

# Expected table data
EXPECTED_HEADERS = ['KPI', 'Q1', 'Q2', 'Q3', 'Q4', 'Full Year']
EXPECTED_DATA = [
    ['Revenue',        '$12M',  '$14.5M', '$16.2M', '$18.8M', '$61.5M'],
    ['Gross Margin',   '58%',   '59%',    '61%',    '62%',    '60%'],
    ['EBITDA',         '$2.1M', '$2.8M',  '$3.5M',  '$4.5M',  '$12.9M'],
    ['Free Cash Flow', '$1.2M', '$1.8M',  '$2.4M',  '$3.2M',  '$8.6M'],
]


def find_table_on_slide(slide):
    """Find first table shape on a slide."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return shape.table
    return None


def get_cell_fill_hex(cell):
    """Extract solid fill color hex from a table cell, or None."""
    tcPr = cell._tc.tcPr
    if tcPr is None:
        return None
    solidFill = tcPr.find(qn('a:solidFill'))
    if solidFill is None:
        return None
    srgb = solidFill.find(qn('a:srgbClr'))
    if srgb is not None:
        return srgb.get('val', '').upper()
    return None


def get_first_run_font(cell):
    """Get font properties of first non-empty run in cell. Returns (name, size_emu, color_hex_or_None)."""
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            name = run.font.name
            size = run.font.size  # EMU or None
            color_hex = None
            try:
                if run.font.color.type is not None:
                    color_hex = str(run.font.color.rgb).upper()
            except Exception:
                pass
            return (name, size, color_hex)
    return (None, None, None)


def verify_task(file_path):
    """Verify task completion with progressive scoring. Returns float 0.0-1.0."""
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # 0-indexed, slide 4
    table = find_table_on_slide(slide)

    # Component 1: Table exists on slide 4 with correct dimensions (0.25 points)
    try:
        if table is None:
            print("FAIL: Component 1 — No table found on slide 4")
        else:
            num_rows = len(table.rows)
            num_cols = len(table.columns)
            if num_rows == 5 and num_cols == 6:
                print(f"PASS: Component 1 — Table found on slide 4 with {num_rows}x{num_cols} (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 — Table dimensions {num_rows}x{num_cols}, expected 5x6")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    if table is None:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Header row text matches (0.25 points)
    try:
        header_match = 0
        for c in range(min(len(table.columns), 6)):
            actual = table.cell(0, c).text.strip()
            expected = EXPECTED_HEADERS[c]
            if actual == expected:
                header_match += 1
            else:
                print(f"  Header mismatch col {c}: expected {repr(expected)}, got {repr(actual)}")

        if header_match == 6:
            print(f"PASS: Component 2 — All 6 header cells match (0.25 pts)")
            total_score += 0.25
        elif header_match > 0:
            partial = round(0.25 * header_match / 6, 3)
            print(f"PARTIAL: Component 2 — {header_match}/6 headers match ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No header cells match")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Data cell values match (0.25 points)
    try:
        data_match = 0
        total_data_cells = 24  # 4 rows x 6 cols
        for r in range(4):
            for c in range(min(len(table.columns), 6)):
                row_idx = r + 1  # skip header
                if row_idx >= len(table.rows):
                    continue
                actual = table.cell(row_idx, c).text.strip()
                expected = EXPECTED_DATA[r][c]
                if actual == expected:
                    data_match += 1
                else:
                    print(f"  Data mismatch [{row_idx},{c}]: expected {repr(expected)}, got {repr(actual)}")

        if data_match == total_data_cells:
            print(f"PASS: Component 3 — All {total_data_cells} data cells match (0.25 pts)")
            total_score += 0.25
        elif data_match > 0:
            partial = round(0.25 * data_match / total_data_cells, 3)
            print(f"PARTIAL: Component 3 — {data_match}/{total_data_cells} data cells match ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No data cells match")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Header row formatting — fill #003366, white text (0.15 points)
    try:
        fill_ok = 0
        text_color_ok = 0
        for c in range(min(len(table.columns), 6)):
            cell = table.cell(0, c)
            fill_hex = get_cell_fill_hex(cell)
            if fill_hex == '003366':
                fill_ok += 1
            else:
                print(f"  Header fill col {c}: expected 003366, got {fill_hex}")

            _, _, color_hex = get_first_run_font(cell)
            if color_hex == 'FFFFFF':
                text_color_ok += 1
            else:
                print(f"  Header text color col {c}: expected FFFFFF, got {color_hex}")

        fill_score = 0.075 * fill_ok / 6
        text_score = 0.075 * text_color_ok / 6
        comp4_score = round(fill_score + text_score, 3)
        if comp4_score >= 0.149:
            print(f"PASS: Component 4 — Header fill ({fill_ok}/6) and text color ({text_color_ok}/6) correct (0.15 pts)")
        elif comp4_score > 0:
            print(f"PARTIAL: Component 4 — fill {fill_ok}/6, text color {text_color_ok}/6 ({comp4_score} pts)")
        else:
            print(f"FAIL: Component 4 — No header formatting matches")
        if comp4_score > 0:
            total_score += comp4_score
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Data cells use 14pt Calibri (0.10 points)
    try:
        font_ok = 0
        total_checked = 0
        for r in range(1, min(len(table.rows), 5)):
            for c in range(min(len(table.columns), 6)):
                total_checked += 1
                name, size, _ = get_first_run_font(table.cell(r, c))
                name_match = (name is not None and name.lower() == 'calibri')
                # 14pt = 177800 EMU
                size_match = (size is not None and abs(size - 177800) < 1000)
                if name_match and size_match:
                    font_ok += 1
                else:
                    if not name_match:
                        print(f"  Font name [{r},{c}]: expected Calibri, got {name}")
                    if not size_match:
                        print(f"  Font size [{r},{c}]: expected 177800 (14pt), got {size}")

        if total_checked > 0 and font_ok == total_checked:
            print(f"PASS: Component 5 — All {total_checked} data cells use 14pt Calibri (0.10 pts)")
            total_score += 0.10
        elif font_ok > 0:
            partial = round(0.10 * font_ok / total_checked, 3)
            print(f"PARTIAL: Component 5 — {font_ok}/{total_checked} data cells correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 — No data cells use 14pt Calibri")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
