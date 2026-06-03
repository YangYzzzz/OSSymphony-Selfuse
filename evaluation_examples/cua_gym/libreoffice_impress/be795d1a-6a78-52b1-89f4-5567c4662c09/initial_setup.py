"""
Initial Setup: Add a financial scorecard table to slide 4
Task ID: impress_exec_056
Domain: libreoffice_impress

Creates an 8-slide financial presentation. Slide 4 has title
'Financial Scorecard' but NO table — the agent's task is to add it.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_056'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def add_title_only_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(prs, "Annual Financial Report",
                    "Fiscal Year 2025 | Prepared by Corporate Finance")

    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "Revenue grew 22% year-over-year to $61.5M driven by enterprise expansion",
        "Gross margins improved from 56% to 60% through operational efficiencies",
        "EBITDA reached $12.9M, a 35% increase over the prior fiscal year",
        "Free cash flow generation of $8.6M supports strategic M&A pipeline",
        "Customer retention rate held steady at 94% across all segments",
    ])

    # Slide 3: Revenue Breakdown
    add_content_slide(prs, "Revenue Breakdown by Segment", [
        "Enterprise Solutions: $28.4M (46% of total revenue)",
        "Mid-Market Accounts: $18.5M (30% of total revenue)",
        "SMB & Self-Service: $9.8M (16% of total revenue)",
        "Professional Services: $4.8M (8% of total revenue)",
        "International markets contributed 23% of overall revenue",
    ])

    # Slide 4: Financial Scorecard — TITLE ONLY, NO TABLE
    add_title_only_slide(prs, "Financial Scorecard")

    # Slide 5: Operational Highlights
    add_content_slide(prs, "Operational Highlights", [
        "Headcount grew from 187 to 234 employees (+25%)",
        "New product launches: 3 major releases, 12 feature updates",
        "Infrastructure costs reduced 18% through cloud optimization",
        "Average deal size increased 15% in enterprise segment",
        "Time-to-close shortened by 8 days for mid-market deals",
    ])

    # Slide 6: Strategic Initiatives
    add_content_slide(prs, "Strategic Initiatives for FY2026", [
        "Launch AI-powered analytics module in Q1 2026",
        "Expand APAC sales presence with Tokyo and Singapore offices",
        "Achieve SOC 2 Type II certification by end of Q2",
        "Target $80M revenue with 63% gross margin",
        "Invest $5M in R&D for next-generation platform capabilities",
    ])

    # Slide 7: Risk Assessment
    add_content_slide(prs, "Risk Assessment", [
        "Currency exposure: 23% revenue in non-USD denominations",
        "Key person dependency in enterprise sales leadership",
        "Regulatory changes in EU data privacy may require platform updates",
        "Competitive pressure from two well-funded Series D entrants",
        "Supply chain constraints could delay hardware partnerships",
    ])

    # Slide 8: Closing & Q&A
    add_title_slide(prs, "Thank You",
                    "Questions & Discussion | investor.relations@acmecorp.com")

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
