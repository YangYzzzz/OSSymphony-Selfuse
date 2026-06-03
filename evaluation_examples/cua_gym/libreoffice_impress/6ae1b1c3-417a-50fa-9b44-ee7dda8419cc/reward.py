"""
Reward Script: Add the title text of each slide as a speaker note on each slide, for all 5 slides.
Task ID: osworld_impress_slide_notes_005
Domain: libreoffice_impress
Scoring:
  Component 1 (0.2): Slide 1 notes == "2025 Product Roadmap Overview"
  Component 2 (0.2): Slide 2 notes == "Q1: Foundation & Infrastructure"
  Component 3 (0.2): Slide 3 notes == "Q2: User Experience Enhancements"
  Component 4 (0.2): Slide 4 notes == "Q3: Growth & Expansion Features"
  Component 5 (0.2): Slide 5 notes == "Q4: Scale & Reliability"
"""

import os

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_slide_notes_005'

# Ground-truth title texts verified from initial_env (idx==0 placeholder)
EXPECTED_NOTES = [
    "2025 Product Roadmap Overview",
    "Q1: Foundation & Infrastructure",
    "Q2: User Experience Enhancements",
    "Q3: Growth & Expansion Features",
    "Q4: Scale & Reliability",
]


def get_slide_title(slide):
    """Return the text of the title placeholder (idx==0) or empty string."""
    for shape in slide.shapes:
        if (hasattr(shape, 'placeholder_format')
                and shape.placeholder_format is not None
                and shape.placeholder_format.idx == 0):
            return shape.text
    return ""


def get_slide_notes(slide):
    """Return stripped notes text, or empty string if none."""
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def verify_task(file_path):
    """
    Verify that each slide's speaker notes contain the slide's title text.
    Returns a float in [0.0, 1.0] with 0.2 per correct slide note.
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have exactly 5 slides
    num_slides = len(prs.slides)
    if num_slides != 5:
        print(f"FAIL: Expected 5 slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    per_slide_score = 0.2

    for slide_idx in range(5):
        slide = prs.slides[slide_idx]
        expected = EXPECTED_NOTES[slide_idx]

        # Component (slide_idx+1): notes of slide == expected title text
        try:
            notes_text = get_slide_notes(slide)
            if notes_text == expected:
                print(f"PASS: Slide {slide_idx + 1} notes == {repr(expected)} ({per_slide_score} pts)")
                total_score += per_slide_score
            else:
                print(f"FAIL: Slide {slide_idx + 1} notes expected {repr(expected)}, found {repr(notes_text)}")
        except Exception as e:
            print(f"ERROR: Slide {slide_idx + 1} notes check failed: {e}")

    final_score = round(min(total_score, 1.0), 4)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
