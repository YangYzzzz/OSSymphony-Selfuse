"""
Initial Setup: 5-slide product roadmap presentation with titles but no speaker notes.
Task ID: osworld_impress_slide_notes_005
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'  # VM path — all scripts run on the VM
TASK_ID = 'osworld_impress_slide_notes_005'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide data: (title, bullets)
    slides_data = [
        (
            "2025 Product Roadmap Overview",
            [
                "Strategic priorities for the year ahead",
                "Aligning engineering, design, and go-to-market",
                "Quarterly milestones and success metrics",
                "Key stakeholder dependencies",
            ],
        ),
        (
            "Q1: Foundation & Infrastructure",
            [
                "Migrate core services to Kubernetes clusters",
                "Launch new CI/CD pipeline with automated testing",
                "Complete data warehouse migration to Snowflake",
                "Reduce average API response time by 30%",
            ],
        ),
        (
            "Q2: User Experience Enhancements",
            [
                "Redesign onboarding flow for enterprise clients",
                "Release mobile app v3.0 for iOS and Android",
                "Introduce AI-powered search and recommendations",
                "Achieve 4.5+ app store rating across platforms",
            ],
        ),
        (
            "Q3: Growth & Expansion Features",
            [
                "Launch multi-tenant SaaS offering in APAC region",
                "Introduce advanced analytics dashboard",
                "Integrate with Salesforce and HubSpot CRM systems",
                "Support SSO via SAML 2.0 and OAuth 2.0",
            ],
        ),
        (
            "Q4: Scale & Reliability",
            [
                "Achieve 99.99% uptime SLA across all tiers",
                "Complete SOC 2 Type II certification",
                "Deploy global CDN for sub-100ms load times",
                "Deliver year-end product retrospective to stakeholders",
            ],
        ),
    ]

    for idx, (title_text, bullets) in enumerate(slides_data):
        if idx == 0:
            # Use Title Slide layout for the first slide
            layout = prs.slide_layouts[0]  # Title Slide
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title_text
            # Subtitle placeholder
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = "Presented by Product Strategy Team · March 2025"
        else:
            # Use Title + Content layout for subsequent slides
            layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title_text
            # Content placeholder
            tf = slide.placeholders[1].text_frame
            tf.text = bullets[0]
            for bullet in bullets[1:]:
                para = tf.add_paragraph()
                para.text = bullet
                para.level = 0

        # IMPORTANT: Do NOT add any speaker notes — the task requires the agent to add them

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
