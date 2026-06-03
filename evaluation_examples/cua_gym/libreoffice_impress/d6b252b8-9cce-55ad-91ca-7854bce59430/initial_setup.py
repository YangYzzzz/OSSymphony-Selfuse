"""
Initial Setup: Before and After comparison on slide 5
Task ID: impress_stu_080
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# --- Create placeholder images first ---
from PIL import Image, ImageDraw, ImageFont

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_080'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
DOWNLOADS = f'{WORKDIR}/Downloads'


def create_placeholder_images():
    """Create realistic before/after treatment images."""
    os.makedirs(DOWNLOADS, exist_ok=True)

    # Before image - a somewhat faded/damaged surface
    img_before = Image.new('RGB', (800, 600), color=(210, 195, 180))
    draw = ImageDraw.Draw(img_before)
    # Add some "damage" marks
    for i in range(0, 800, 60):
        for j in range(0, 600, 80):
            draw.ellipse([i, j, i + 25, j + 25], fill=(180, 160, 140))
    draw.rectangle([50, 50, 750, 550], outline=(150, 130, 110), width=3)
    draw.text((300, 280), "BEFORE", fill=(120, 100, 80))
    img_before.save(f'{DOWNLOADS}/before.png')

    # After image - a clean, restored surface
    img_after = Image.new('RGB', (800, 600), color=(240, 245, 250))
    draw = ImageDraw.Draw(img_after)
    draw.rectangle([50, 50, 750, 550], outline=(100, 140, 200), width=3)
    # Add subtle highlight areas
    for i in range(100, 700, 120):
        draw.rectangle([i, 100, i + 80, 500], fill=(235, 240, 248))
    draw.text((320, 280), "AFTER", fill=(60, 100, 180))
    img_after.save(f'{DOWNLOADS}/after.png')

    print(f'Created before.png and after.png in {DOWNLOADS}')


def create_initial():
    prs = Presentation()
    # Standard 16:9 slide size (10x7.5)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Treatment Results"
    slide1.placeholders[1].text = "Dermatological Treatment Efficacy Study\nQ1 2025 Report"

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Study Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Study Duration: January 2024 - December 2024"
    p = body2.add_paragraph()
    p.text = "Total Participants: 245 patients across 3 clinical sites"
    p = body2.add_paragraph()
    p.text = "Treatment Protocol: Combination therapy with topical retinoid and laser resurfacing"
    p = body2.add_paragraph()
    p.text = "Follow-up Period: 6 months post-treatment"

    # --- Slide 3: Methodology ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Methodology"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Phase 1: Baseline assessment and skin biopsy"
    p = body3.add_paragraph()
    p.text = "Phase 2: Initial treatment cycle (8 weeks)"
    p = body3.add_paragraph()
    p.text = "Phase 3: Maintenance therapy (16 weeks)"
    p = body3.add_paragraph()
    p.text = "Phase 4: Post-treatment evaluation and photography"

    # --- Slide 4: Key Findings ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Key Findings"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "87% of patients showed significant improvement"
    p = body4.add_paragraph()
    p.text = "Average skin texture score improved from 3.2 to 7.8 (out of 10)"
    p = body4.add_paragraph()
    p.text = "Hyperpigmentation reduced by 64% on average"
    p = body4.add_paragraph()
    p.text = "Patient satisfaction rate: 92%"

    # --- Slide 5: Visual Comparison (EMPTY - agent's workspace) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add only a title textbox at the top
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Visual Comparison"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    # NO other shapes on slide 5 - agent must add divider, images, and labels

    # --- Slide 6: Statistical Analysis ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Statistical Analysis"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "p-value < 0.001 for primary endpoint"
    p = body6.add_paragraph()
    p.text = "95% CI for improvement: [62.3%, 71.8%]"
    p = body6.add_paragraph()
    p.text = "Effect size (Cohen's d): 1.42 (large effect)"
    p = body6.add_paragraph()
    p.text = "No significant difference between clinical sites (p = 0.73)"

    # --- Slide 7: Adverse Events ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Adverse Events"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "Mild erythema: 23% of patients (resolved within 48 hours)"
    p = body7.add_paragraph()
    p.text = "Temporary sensitivity: 15% of patients"
    p = body7.add_paragraph()
    p.text = "No serious adverse events reported"
    p = body7.add_paragraph()
    p.text = "Treatment discontinuation rate: 3.2%"

    # --- Slide 8: Conclusions ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Conclusions & Next Steps"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Combination therapy demonstrates superior efficacy"
    p = body8.add_paragraph()
    p.text = "Recommend expansion to Phase III multi-center trial"
    p = body8.add_paragraph()
    p.text = "Long-term follow-up study planned for 2025-2026"
    p = body8.add_paragraph()
    p.text = "Publication target: Journal of Dermatological Treatment"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


create_placeholder_images()
create_initial()

# GUI-ready startup
launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')
