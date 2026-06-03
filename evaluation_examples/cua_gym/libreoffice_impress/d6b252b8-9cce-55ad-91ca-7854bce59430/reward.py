"""
Reward Script: Before/After comparison on slide 5
Task ID: impress_stu_080
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Vertical center divider line on slide 5
  Component 2 (0.25): Before image on left side, ~4 inches wide
  Component 3 (0.25): After image on right side, ~4 inches wide
  Component 4 (0.125): 'Before Treatment' label at 16pt
  Component 5 (0.125): 'After Treatment' label at 16pt
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_080'
FILE_PATH = f'{WORKDIR}/{TASK_ID}.pptx'

# Tolerances
INCH_TOLERANCE = 0.3  # inches tolerance for position/size checks
PT_TOLERANCE = 2  # pt tolerance for font size


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # slide 5 (0-indexed)
    slide_width = prs.slide_width
    slide_center_x = slide_width // 2  # center of slide in EMU

    # Collect shapes by type for analysis
    lines = []
    pictures = []
    textboxes = []

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.LINE or shape.shape_type == 9:
            lines.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.shape_type == 13:
            pictures.append(shape)
        elif hasattr(shape, 'text_frame') and shape.text and shape.shape_type != 14:
            # Exclude placeholder shapes that are pre-existing (type 14 = PLACEHOLDER)
            textboxes.append(shape)

    # =========================================================================
    # Component 1: Vertical center divider line (0.25 points)
    # A vertical line near the center of the slide
    # =========================================================================
    try:
        found_center_line = False
        for line in lines:
            # A vertical line has width ~0 and significant height
            is_vertical = (line.width <= Inches(0.2)) and (line.height >= Inches(1.0))
            # Check if positioned near center
            line_x = line.left + line.width // 2
            near_center = abs(line_x - slide_center_x) <= Inches(1.0)

            if is_vertical and near_center:
                found_center_line = True
                print(f"PASS: Component 1 — Vertical center divider found (name={line.name}, x={line.left}, h={line.height}) (0.25 pts)")
                total_score += 0.25
                break

        if not found_center_line:
            print(f"FAIL: Component 1 — No vertical center divider line found on slide 5. Lines found: {len(lines)}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # =========================================================================
    # Component 2: Before image on left side, ~4 inches wide (0.25 points)
    # Must be an image on the LEFT half of the slide, approximately 4 inches wide,
    # and its blob must match before.png
    # =========================================================================
    try:
        before_blob = open('/home/user/Downloads/before.png', 'rb').read()
        found_before_img = False

        for pic in pictures:
            pic_center_x = pic.left + pic.width // 2
            is_left = pic_center_x < slide_center_x
            is_before = (pic.image.blob == before_blob)
            width_inches = pic.width / 914400.0
            width_ok = abs(width_inches - 4.0) <= INCH_TOLERANCE

            if is_before and is_left and width_ok:
                found_before_img = True
                print(f"PASS: Component 2 — Before image on left side, width={width_inches:.2f}in (0.25 pts)")
                total_score += 0.25
                break

        if not found_before_img:
            # Check partial: any image matching before.png?
            any_before = any(pic.image.blob == before_blob for pic in pictures)
            if any_before:
                print(f"FAIL: Component 2 — before.png found but not on left side or wrong width")
            else:
                print(f"FAIL: Component 2 — before.png not found on slide 5. Pictures: {len(pictures)}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # =========================================================================
    # Component 3: After image on right side, ~4 inches wide (0.25 points)
    # =========================================================================
    try:
        after_blob = open('/home/user/Downloads/after.png', 'rb').read()
        found_after_img = False

        for pic in pictures:
            pic_center_x = pic.left + pic.width // 2
            is_right = pic_center_x > slide_center_x
            is_after = (pic.image.blob == after_blob)
            width_inches = pic.width / 914400.0
            width_ok = abs(width_inches - 4.0) <= INCH_TOLERANCE

            if is_after and is_right and width_ok:
                found_after_img = True
                print(f"PASS: Component 3 — After image on right side, width={width_inches:.2f}in (0.25 pts)")
                total_score += 0.25
                break

        if not found_after_img:
            any_after = any(pic.image.blob == after_blob for pic in pictures)
            if any_after:
                print(f"FAIL: Component 3 — after.png found but not on right side or wrong width")
            else:
                print(f"FAIL: Component 3 — after.png not found on slide 5. Pictures: {len(pictures)}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # =========================================================================
    # Component 4: 'Before Treatment' label at 16pt (0.125 points)
    # A text shape containing 'Before Treatment' with font size ~16pt
    # =========================================================================
    try:
        found_before_label = False
        target_size_emu = Pt(16)  # 203200 EMU

        for tb in textboxes:
            text_lower = tb.text.strip().lower()
            if 'before treatment' in text_lower:
                # Check font size
                for para in tb.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size is not None:
                            size_pt = run.font.size / 12700.0
                            if abs(size_pt - 16.0) <= PT_TOLERANCE:
                                found_before_label = True
                                print(f"PASS: Component 4 — 'Before Treatment' label at {size_pt:.1f}pt (0.125 pts)")
                                total_score += 0.125
                                break
                    if found_before_label:
                        break
            if found_before_label:
                break

        if not found_before_label:
            # Check if text exists at all
            has_text = any('before treatment' in tb.text.strip().lower() for tb in textboxes)
            if has_text:
                print(f"FAIL: Component 4 — 'Before Treatment' text found but font size not ~16pt")
            else:
                print(f"FAIL: Component 4 — 'Before Treatment' label not found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # =========================================================================
    # Component 5: 'After Treatment' label at 16pt (0.125 points)
    # =========================================================================
    try:
        found_after_label = False

        for tb in textboxes:
            text_lower = tb.text.strip().lower()
            if 'after treatment' in text_lower:
                for para in tb.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size is not None:
                            size_pt = run.font.size / 12700.0
                            if abs(size_pt - 16.0) <= PT_TOLERANCE:
                                found_after_label = True
                                print(f"PASS: Component 5 — 'After Treatment' label at {size_pt:.1f}pt (0.125 pts)")
                                total_score += 0.125
                                break
                    if found_after_label:
                        break
            if found_after_label:
                break

        if not found_after_label:
            has_text = any('after treatment' in tb.text.strip().lower() for tb in textboxes)
            if has_text:
                print(f"FAIL: Component 5 — 'After Treatment' text found but font size not ~16pt")
            else:
                print(f"FAIL: Component 5 — 'After Treatment' label not found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(FILE_PATH)
