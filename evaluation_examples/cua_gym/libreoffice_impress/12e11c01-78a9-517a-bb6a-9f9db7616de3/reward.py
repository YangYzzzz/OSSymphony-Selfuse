"""
Reward Script: Insert dna_structure.png image onto slide 3, right half, ~4x5 inches
Task ID: impress_stu_018
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): Slide 3 contains at least one picture shape
  Component 2 (0.3): The picture is the correct image (blob matches dna_structure.png)
  Component 3 (0.2): Image positioned on right half of the slide
  Component 4 (0.2): Image sized approximately 4 inches wide x 5 inches tall
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_018'
SOURCE_IMAGE_PATH = '/home/user/Downloads/dna_structure.png'

# Tolerance for size check: 15% relative tolerance
SIZE_TOLERANCE = 0.15

# Target dimensions in EMU (914400 EMU = 1 inch)
TARGET_WIDTH_EMU = Inches(4)   # 3657600
TARGET_HEIGHT_EMU = Inches(5)  # 4572000


def is_approx_equal(actual, expected, tolerance=SIZE_TOLERANCE):
    """Check if actual is within tolerance of expected (relative)."""
    if expected == 0:
        return actual == 0
    return abs(actual - expected) / expected <= tolerance


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Presentation has at least 3 slides
    if len(prs.slides) < 3:
        print(f"PRECONDITION FAIL: Need at least 3 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]  # 0-indexed, slide 3
    slide_width = prs.slide_width

    # Find all picture shapes on slide 3
    pictures = [s for s in slide3.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]

    # Load source image blob for comparison
    source_blob = None
    try:
        with open(SOURCE_IMAGE_PATH, 'rb') as f:
            source_blob = f.read()
    except Exception as e:
        print(f"WARNING: Cannot read source image for comparison: {e}")

    # Component 1: Slide 3 contains at least one picture shape (0.3 points)
    # This FAILS on initial (no pictures) -> PASSES on golden (has picture)
    try:
        if len(pictures) > 0:
            print(f"PASS: Component 1 - Slide 3 has {len(pictures)} picture(s) (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 - Slide 3 has no pictures")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: The picture is the correct image matching dna_structure.png (0.3 points)
    # This FAILS on initial (no pictures) -> PASSES on golden (correct blob)
    matching_pic = None
    try:
        if source_blob is not None and len(pictures) > 0:
            for pic in pictures:
                if pic.image.blob == source_blob:
                    matching_pic = pic
                    break
            if matching_pic is not None:
                print(f"PASS: Component 2 - Image blob matches dna_structure.png ({len(source_blob)} bytes) (0.3 pts)")
                total_score += 0.3
            else:
                blob_sizes = [len(p.image.blob) for p in pictures]
                print(f"FAIL: Component 2 - No picture blob matches source. Picture blob sizes: {blob_sizes}, source: {len(source_blob)}")
        elif len(pictures) == 0:
            print(f"FAIL: Component 2 - No pictures on slide 3 to check")
        else:
            print(f"FAIL: Component 2 - Could not load source image for comparison")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Use matching_pic if found, otherwise use first picture for position/size checks
    target_pic = matching_pic if matching_pic else (pictures[0] if pictures else None)

    # Component 3: Image is positioned on the right half of the slide (0.2 points)
    # Right half means the image's left edge is at or past the midpoint of the slide
    # This FAILS on initial (no pictures) -> PASSES on golden (positioned right)
    try:
        if target_pic is not None:
            midpoint = slide_width / 2
            pic_left = target_pic.left
            if pic_left >= midpoint:
                print(f"PASS: Component 3 - Image left={pic_left/914400:.2f}in >= midpoint={midpoint/914400:.2f}in (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 3 - Image left={pic_left/914400:.2f}in < midpoint={midpoint/914400:.2f}in. Not on right half.")
        else:
            print(f"FAIL: Component 3 - No picture to check position")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Image is approximately 4 inches wide x 5 inches tall (0.2 points)
    # This FAILS on initial (no pictures) -> PASSES on golden (correct size)
    try:
        if target_pic is not None:
            actual_w = target_pic.width
            actual_h = target_pic.height
            w_ok = is_approx_equal(actual_w, TARGET_WIDTH_EMU)
            h_ok = is_approx_equal(actual_h, TARGET_HEIGHT_EMU)
            if w_ok and h_ok:
                print(f"PASS: Component 4 - Size {actual_w/914400:.2f}x{actual_h/914400:.2f}in ~= 4.00x5.00in (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 4 - Size {actual_w/914400:.2f}x{actual_h/914400:.2f}in, expected ~4.00x5.00in (w_ok={w_ok}, h_ok={h_ok})")
        else:
            print(f"FAIL: Component 4 - No picture to check size")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
