"""
Initial Setup: Insert an image onto slide 3 of a Genetics Lecture presentation.
Task ID: impress_stu_018
Domain: libreoffice_impress

Creates a 6-slide Genetics Lecture presentation with text content.
Slide 3 has title 'DNA Structure' and text on the left side, no image.
Also creates /home/user/Downloads/dna_structure.png as source image.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a simple DNA structure image using PIL
from PIL import Image, ImageDraw, ImageFont

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_018'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
IMG_DIR = f'{WORKDIR}/Downloads'
IMG_PATH = f'{IMG_DIR}/dna_structure.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_dna_image():
    """Create a realistic-looking DNA structure diagram image."""
    os.makedirs(IMG_DIR, exist_ok=True)

    width, height = 600, 750
    img = Image.new('RGB', (width, height), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)

    # Draw a stylized double helix
    import math
    center_x = width // 2
    amplitude = 100
    for y in range(20, height - 20, 2):
        phase = y * 0.05
        x1 = center_x + int(amplitude * math.sin(phase))
        x2 = center_x + int(amplitude * math.sin(phase + math.pi))

        # Left strand - blue
        draw.ellipse([x1 - 4, y - 2, x1 + 4, y + 2], fill=(0, 100, 200))
        # Right strand - red
        draw.ellipse([x2 - 4, y - 2, x2 + 4, y + 2], fill=(200, 50, 50))

        # Base pair connections every 30 pixels
        if y % 30 < 3:
            colors = [(50, 180, 50), (200, 150, 0), (150, 50, 200), (200, 100, 50)]
            color = colors[(y // 30) % len(colors)]
            draw.line([x1, y, x2, y], fill=color, width=3)

    # Add label
    draw.text((width // 2 - 80, height - 40), "DNA Double Helix", fill=(0, 0, 0))

    img.save(IMG_PATH)
    print(f'DNA image created: {IMG_PATH}')


def create_initial():
    prs = Presentation()
    # Standard widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Genetics Lecture"
    slide1.placeholders[1].text = "Introduction to Molecular Biology\nDr. Elena Rodriguez\nSpring 2026"

    # --- Slide 2: Introduction to Genetics ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Introduction to Genetics"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Genetics is the study of heredity and variation in living organisms"
    p2a = body2.add_paragraph()
    p2a.text = "Key concepts include DNA replication, transcription, and translation"
    p2a.level = 1
    p2b = body2.add_paragraph()
    p2b.text = "Modern genetics encompasses genomics, epigenetics, and gene therapy"
    p2b.level = 1
    p2c = body2.add_paragraph()
    p2c.text = "Understanding genetics is fundamental to medicine, agriculture, and forensics"
    p2c.level = 1

    # --- Slide 3: DNA Structure (text on left side, NO image) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Title at top
    title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(5), Inches(1))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "DNA Structure"
    run_title = p_title.runs[0]
    run_title.font.size = Pt(36)
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Content text on the LEFT side
    content_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    lines = [
        ("The DNA molecule consists of two polynucleotide chains", Pt(18), False),
        ("", Pt(14), False),
        ("Key Components:", Pt(20), True),
        ("  - Deoxyribose sugar backbone", Pt(16), False),
        ("  - Phosphate groups linking nucleotides", Pt(16), False),
        ("  - Nitrogenous bases: Adenine, Thymine, Guanine, Cytosine", Pt(16), False),
        ("  - Hydrogen bonds between complementary base pairs", Pt(16), False),
        ("", Pt(14), False),
        ("Watson and Crick proposed the double helix model in 1953,", Pt(16), False),
        ("building on X-ray crystallography data from Rosalind Franklin.", Pt(16), False),
        ("", Pt(14), False),
        ("The antiparallel arrangement of the two strands runs", Pt(16), False),
        ("5' to 3' in opposite directions.", Pt(16), False),
    ]

    p0 = tf.paragraphs[0]
    p0.text = lines[0][0]
    r0 = p0.runs[0]
    r0.font.size = lines[0][1]

    for text, size, bold in lines[1:]:
        p = tf.add_paragraph()
        p.text = text
        if p.runs:
            p.runs[0].font.size = size
            p.runs[0].font.bold = bold
            p.runs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 4: Gene Expression ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Gene Expression"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Gene expression is the process by which DNA directs protein synthesis"
    p4a = body4.add_paragraph()
    p4a.text = "Transcription: DNA is copied to messenger RNA (mRNA) in the nucleus"
    p4a.level = 1
    p4b = body4.add_paragraph()
    p4b.text = "Translation: Ribosomes read mRNA codons to assemble amino acid chains"
    p4b.level = 1
    p4c = body4.add_paragraph()
    p4c.text = "Post-translational modifications fold proteins into functional shapes"
    p4c.level = 1
    p4d = body4.add_paragraph()
    p4d.text = "Regulation occurs at multiple levels: chromatin, transcription, and epigenetic"
    p4d.level = 1

    # --- Slide 5: Mutations and Variations ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Mutations and Variations"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Mutations are changes in the DNA sequence that may affect phenotype"
    p5a = body5.add_paragraph()
    p5a.text = "Point mutations: substitution of a single nucleotide base pair"
    p5a.level = 1
    p5b = body5.add_paragraph()
    p5b.text = "Frameshift mutations: insertions or deletions altering the reading frame"
    p5b.level = 1
    p5c = body5.add_paragraph()
    p5c.text = "Chromosomal abnormalities: large-scale structural or numerical changes"
    p5c.level = 1
    p5d = body5.add_paragraph()
    p5d.text = "SNPs (Single Nucleotide Polymorphisms) account for most human genetic variation"
    p5d.level = 1

    # --- Slide 6: Summary ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Summary"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "DNA stores genetic information in a double helix structure"
    p6a = body6.add_paragraph()
    p6a.text = "Gene expression converts DNA instructions into functional proteins"
    p6a.level = 1
    p6b = body6.add_paragraph()
    p6b.text = "Mutations introduce genetic variation essential for evolution"
    p6b.level = 1
    p6c = body6.add_paragraph()
    p6c.text = "Modern tools like CRISPR-Cas9 enable precise genome editing"
    p6c.level = 1
    p6d = body6.add_paragraph()
    p6d.text = "Next lecture: Mendelian Inheritance and Punnett Squares"
    p6d.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_dna_image()
create_initial()
