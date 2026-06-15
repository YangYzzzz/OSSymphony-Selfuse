"""
Initial Setup: Set body text on slide 5 to exactly 24pt line spacing
Task ID: impress_tct_089
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_089'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_paragraph(tf, text, font_size=Pt(18), bold=False, alignment=None, color=None):
    """Add a paragraph with specified formatting to a text frame."""
    p = tf.add_paragraph()
    p.text = text
    if alignment:
        p.alignment = alignment
    for run in p.runs:
        run.font.size = font_size
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
    return p


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Precise Layout Design"
    slide1.placeholders[1].text = "Typography & Spacing Workshop 2025"

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Workshop Objectives"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Understand the fundamentals of precise text layout in presentations"
    for line in [
        "Master the difference between proportional and fixed line spacing",
        "Learn when to use exact spacing for consistent visual rhythm",
        "Apply spacing techniques to real-world slide designs",
        "Evaluate readability across different font sizes and weights",
    ]:
        add_text_paragraph(body2, line, font_size=Pt(16))

    # --- Slide 3: Spacing Types ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Types of Line Spacing"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Proportional Spacing (Single, 1.5, Double)"
    for line in [
        "Adapts automatically based on the largest font in each line",
        "Good default choice for mixed content slides",
        "Fixed/Exact Spacing (e.g., 24pt, 30pt)",
        "Maintains constant distance regardless of font size",
        "Ideal for precise layout control and uniform appearance",
    ]:
        add_text_paragraph(body3, line, font_size=Pt(16))

    # --- Slide 4: Design Principles ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Layout Design Principles"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Consistency in vertical rhythm improves readability"
    for line in [
        "Align text baselines across columns when using multi-column layouts",
        "Use exact spacing when mixing font sizes within a text block",
        "Leave sufficient margins between text blocks and slide edges",
        "Consider audience distance when selecting spacing values",
    ]:
        add_text_paragraph(body4, line, font_size=Pt(16))

    # --- Slide 5: Case Study (TARGET SLIDE) ---
    # This slide has 8 lines of body text with DEFAULT proportional spacing
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    slide5.shapes.title_text = None

    # Add title text box
    title_box = slide5.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(1.0))
    tf_title = title_box.text_frame
    tf_title.paragraphs[0].text = "Case Study: Annual Report Cover Page"
    for run in tf_title.paragraphs[0].runs:
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Add body text box with 8 lines — DEFAULT single (proportional) spacing
    body_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.0))
    tf_body = body_box.text_frame
    tf_body.word_wrap = True

    body_lines = [
        "The annual report cover page requires precise vertical alignment of all text elements.",
        "Each line must maintain consistent spacing to create a professional, polished appearance.",
        "When using proportional spacing, lines with taller characters receive extra vertical room.",
        "This inconsistency becomes visible when mixing headings and body text on the same page.",
        "Fixed spacing at 24 points ensures every line occupies exactly the same vertical distance.",
        "Designers at Meridian Publishing adopted this technique for their 2024 corporate materials.",
        "Client feedback indicated a 40% improvement in perceived document quality after the change.",
        "We recommend applying exact spacing to all body text blocks in formal presentation layouts.",
    ]

    tf_body.paragraphs[0].text = body_lines[0]
    for run in tf_body.paragraphs[0].runs:
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    for line in body_lines[1:]:
        p = tf_body.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 6: Summary ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Key Takeaways"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Exact line spacing provides consistent vertical rhythm across all text"
    for line in [
        "Use 24pt exact spacing for standard body text in formal presentations",
        "Always preview slides at actual display resolution to verify spacing",
        "Document your spacing choices in the style guide for team consistency",
        "Test with different projectors — spacing perception varies by display",
    ]:
        add_text_paragraph(body6, line, font_size=Pt(16))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
