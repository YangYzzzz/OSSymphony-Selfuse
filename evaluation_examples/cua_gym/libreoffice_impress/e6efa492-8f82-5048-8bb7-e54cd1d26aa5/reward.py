"""
Reward Script: Set body text on slide 5 to exactly 24pt fixed line spacing
Task ID: impress_tct_089
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Content text box uses fixed (spcPts) line spacing type
  Component 2 (0.3): Spacing value is exactly 24pt (2400 hundredths of a point)
  Component 3 (0.3): All paragraphs have consistent 24pt fixed spacing
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_089'


def persist_app_state(domain):
    """Save any unsaved GUI edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.oxml.ns import qn
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 5 slides
    if len(prs.slides) < 5:
        print(f"CRITICAL: Expected at least 5 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # slide 5 (0-indexed)

    # Find the content text box (the one with multiple paragraphs of body text)
    content_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and len(shape.text_frame.paragraphs) >= 4:
            content_shape = shape
            break

    if content_shape is None:
        print("CRITICAL: No content text box with multiple paragraphs found on slide 5")
        print("REWARD: 0.0")
        return 0.0

    tf = content_shape.text_frame
    paragraphs = tf.paragraphs
    num_paragraphs = len(paragraphs)
    print(f"INFO: Found content text box '{content_shape.name}' with {num_paragraphs} paragraphs")

    # Helper: extract line spacing info from a paragraph's XML
    def get_line_spacing_info(para):
        """Returns (spacing_type, spacing_val) where:
           spacing_type: 'spcPts' (fixed), 'spcPct' (proportional), or None (default)
           spacing_val: integer value attribute, or None
        """
        pPr = para._p.find(qn('a:pPr'))
        if pPr is None:
            return (None, None)
        lnSpc = pPr.find(qn('a:lnSpc'))
        if lnSpc is None:
            return (None, None)
        spcPts = lnSpc.find(qn('a:spcPts'))
        if spcPts is not None:
            return ('spcPts', int(spcPts.get('val', '0')))
        spcPct = lnSpc.find(qn('a:spcPct'))
        if spcPct is not None:
            return ('spcPct', int(spcPct.get('val', '0')))
        return (None, None)

    # Component 1: At least one paragraph uses fixed (spcPts) line spacing (0.4 points)
    # This checks that the spacing TYPE changed from default/proportional to fixed.
    # On initial_env, all paragraphs have None (default) spacing, so this will FAIL.
    try:
        fixed_count = 0
        for para in paragraphs:
            spc_type, spc_val = get_line_spacing_info(para)
            if spc_type == 'spcPts':
                fixed_count += 1

        if fixed_count > 0:
            print(f"PASS: Component 1 -- {fixed_count}/{num_paragraphs} paragraphs use fixed (spcPts) spacing (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 -- No paragraphs use fixed (spcPts) spacing; all use default/proportional")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: The fixed spacing value is exactly 24pt (2400 hundredths of a point) (0.3 points)
    # spcPts val is in hundredths of a point, so 24pt = 2400
    try:
        correct_value_count = 0
        for para in paragraphs:
            spc_type, spc_val = get_line_spacing_info(para)
            if spc_type == 'spcPts' and spc_val == 2400:
                correct_value_count += 1

        if correct_value_count > 0:
            print(f"PASS: Component 2 -- {correct_value_count}/{num_paragraphs} paragraphs have exactly 24pt (2400) spacing (0.3 pts)")
            total_score += 0.3
        else:
            # Check what values exist for debugging
            values_found = []
            for para in paragraphs:
                spc_type, spc_val = get_line_spacing_info(para)
                if spc_type is not None:
                    values_found.append(f"{spc_type}={spc_val}")
            print(f"FAIL: Component 2 -- No paragraphs have 24pt fixed spacing. Values found: {values_found}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: ALL paragraphs consistently have 24pt fixed spacing (0.3 points)
    # This ensures the change was applied to all text, not just some paragraphs.
    try:
        all_correct = True
        for j, para in enumerate(paragraphs):
            spc_type, spc_val = get_line_spacing_info(para)
            if spc_type != 'spcPts' or spc_val != 2400:
                all_correct = False
                print(f"  INFO: Para {j} has spacing type={spc_type}, val={spc_val} (expected spcPts=2400)")

        if all_correct:
            print(f"PASS: Component 3 -- All {num_paragraphs} paragraphs have consistent 24pt fixed spacing (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 3 -- Not all paragraphs have 24pt fixed spacing")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
