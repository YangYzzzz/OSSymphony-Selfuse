"""
Reward Script: Laptop Comparison Presentation (8 slides)
Task ID: impress_wf_082
Domain: libreoffice_impress
Scoring:
  C1: Slide count == 8                          (0.15)
  C2: ECEFF1 backgrounds on all 8 slides        (0.10)
  C3: Wipe transitions on all 8 slides          (0.10)
  C4: Slide 3 has spec comparison table (8x4)   (0.15)
  C5: Slide 5 has rounded rectangle gauges      (0.15)
  C6: Slide 7 has green + red colored text      (0.15)
  C7: Slide 8 has gold (#FFB300) border shape   (0.10)
  C8: Blue (#1565C0) accent used                (0.10)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_082'


def persist_app_state(domain: str):
    """Best-effort save via Ctrl+S before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            import time
            time.sleep(1.0)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load presentation
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)

    # ---------------------------------------------------------------
    # Component 1: Slide count == 8  (0.15 points)
    # Initial has 1 slide; golden has 8.
    # ---------------------------------------------------------------
    try:
        n_slides = len(slides)
        if n_slides == 8:
            print(f"PASS: C1 — Slide count is 8 (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: C1 — Expected 8 slides, found {n_slides}")
    except Exception as e:
        print(f"ERROR: C1 — {e}")

    # Gate: need at least 8 slides for remaining checks
    if len(slides) < 8:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # ---------------------------------------------------------------
    # Component 2: All 8 slides have ECEFF1 solid background (0.10 pts)
    # Initial has inherited/theme background (type 5), golden has solid ECEFF1.
    # ---------------------------------------------------------------
    try:
        bg_pass_count = 0
        for i, slide in enumerate(slides):
            fill = slide.background.fill
            if fill.type == 1:  # SOLID
                try:
                    rgb = str(fill.fore_color.rgb)
                    if rgb.upper() == "ECEFF1":
                        bg_pass_count += 1
                except:
                    pass
        if bg_pass_count == 8:
            print(f"PASS: C2 — All 8 slides have #ECEFF1 background (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: C2 — {bg_pass_count}/8 slides have #ECEFF1 background")
    except Exception as e:
        print(f"ERROR: C2 — {e}")

    # ---------------------------------------------------------------
    # Component 3: Wipe transitions on all 8 slides (0.10 pts)
    # Initial has no transitions; golden has wipe on all.
    # ---------------------------------------------------------------
    try:
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        wipe_count = 0
        with zipfile.ZipFile(file_path, 'r') as zf:
            for i in range(1, 9):
                fname = f'ppt/slides/slide{i}.xml'
                try:
                    with zf.open(fname) as f:
                        root = ET.parse(f).getroot()
                        tr = root.find('.//p:transition', ns)
                        if tr is not None and tr.find(f'.//p:wipe', ns) is not None:
                            wipe_count += 1
                except KeyError:
                    pass
        if wipe_count == 8:
            print(f"PASS: C3 — All 8 slides have Wipe transitions (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: C3 — {wipe_count}/8 slides have Wipe transitions")
    except Exception as e:
        print(f"ERROR: C3 — {e}")

    # ---------------------------------------------------------------
    # Component 4: Slide 3 has a comparison table with 8 rows x 4 cols (0.15 pts)
    # Initial has no table; golden has spec table on slide 3.
    # ---------------------------------------------------------------
    try:
        slide3 = slides[2]
        table_found = False
        for shape in slide3.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                t = shape.table
                n_rows = len(t.rows)
                n_cols = len(t.columns)
                # Expect 8 rows (header + 7 specs) and 4 cols (Spec + 3 models)
                if n_rows >= 7 and n_cols >= 4:
                    # Check header row has spec categories
                    header_text = [t.cell(0, c).text.strip().lower() for c in range(n_cols)]
                    has_spec_header = any('spec' in h or 'model' in h or 'probook' in h or 'techair' in h or 'ultraslim' in h for h in header_text)
                    # Check at least some data rows mention CPU/RAM/Price
                    all_cell_text = ' '.join(t.cell(r, 0).text for r in range(n_rows)).lower()
                    has_cpu = 'cpu' in all_cell_text
                    has_ram = 'ram' in all_cell_text
                    has_price = 'price' in all_cell_text

                    if has_spec_header and has_cpu and has_ram and has_price:
                        print(f"PASS: C4 — Slide 3 table {n_rows}x{n_cols} with CPU/RAM/Price specs (0.15 pts)")
                        total_score += 0.15
                        table_found = True
                    else:
                        print(f"FAIL: C4 — Table found {n_rows}x{n_cols} but missing expected content (header={has_spec_header}, cpu={has_cpu}, ram={has_ram}, price={has_price})")
                        table_found = True
                else:
                    print(f"FAIL: C4 — Table found but wrong size: {n_rows}x{n_cols} (need >=7x4)")
                    table_found = True
                break
        if not table_found:
            print("FAIL: C4 — No table found on slide 3")
    except Exception as e:
        print(f"ERROR: C4 — {e}")

    # ---------------------------------------------------------------
    # Component 5: Slide 5 has rounded rectangle shapes (battery gauges) (0.15 pts)
    # Initial has no such shapes; golden has 6 rounded rectangles on slide 5.
    # ---------------------------------------------------------------
    try:
        slide5 = slides[4]
        rounded_rect_count = 0
        for shape in slide5.shapes:
            if 'Rounded Rectangle' in shape.name or 'Rounded' in shape.name:
                rounded_rect_count += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # Also check via XML for rounded rectangle auto shape type
                try:
                    sp = shape._element
                    prstGeom = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
                    if prstGeom is not None and prstGeom.get('prst') == 'roundRect':
                        rounded_rect_count += 1
                except:
                    pass
        # We expect at least 3 pairs (background + fill) = 6, but at minimum 3
        if rounded_rect_count >= 3:
            print(f"PASS: C5 — Slide 5 has {rounded_rect_count} rounded rectangles (battery gauges) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: C5 — Expected >=3 rounded rectangles on slide 5, found {rounded_rect_count}")
    except Exception as e:
        print(f"ERROR: C5 — {e}")

    # ---------------------------------------------------------------
    # Component 6: Slide 7 has green and red colored text (pros/cons) (0.15 pts)
    # Initial has no colored text; golden has green (2E7D32) pros and red (C62828) cons.
    # ---------------------------------------------------------------
    try:
        slide7 = slides[6]
        has_green = False
        has_red = False
        for shape in slide7.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color.type is not None:
                                rgb = str(run.font.color.rgb).upper()
                                # Green tones: check for greenish colors
                                r_val = int(rgb[0:2], 16)
                                g_val = int(rgb[2:4], 16)
                                b_val = int(rgb[4:6], 16)
                                if g_val > r_val and g_val > b_val and g_val > 80:
                                    has_green = True
                                # Red tones: check for reddish colors
                                if r_val > g_val and r_val > b_val and r_val > 80:
                                    has_red = True
                        except:
                            pass
        if has_green and has_red:
            print(f"PASS: C6 — Slide 7 has green (pros) and red (cons) text (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: C6 — Slide 7 green={has_green}, red={has_red}")
    except Exception as e:
        print(f"ERROR: C6 — {e}")

    # ---------------------------------------------------------------
    # Component 7: Slide 8 has a shape with gold (#FFB300) border (0.10 pts)
    # Initial has no gold border; golden has a winner shape with FFB300 line.
    # ---------------------------------------------------------------
    try:
        slide8 = slides[7]
        gold_border_found = False
        for shape in slide8.shapes:
            try:
                # Check via python-pptx API
                line_rgb = str(shape.line.color.rgb).upper()
                if line_rgb in ('FFB300', 'FFC107', 'FFD700', 'FFAB00'):
                    gold_border_found = True
                    break
                # Also accept close gold tones
                r_val = int(line_rgb[0:2], 16)
                g_val = int(line_rgb[2:4], 16)
                b_val = int(line_rgb[4:6], 16)
                if r_val > 200 and g_val > 150 and b_val < 80:
                    gold_border_found = True
                    break
            except:
                pass
            # Fallback: check XML directly
            try:
                ln = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
                if ln is not None:
                    sf = ln.find('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
                    if sf is not None:
                        srgb = sf.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                        if srgb is not None:
                            val = srgb.get('val', '').upper()
                            if val in ('FFB300', 'FFC107', 'FFD700', 'FFAB00'):
                                gold_border_found = True
                                break
                            r_val = int(val[0:2], 16)
                            g_val = int(val[2:4], 16)
                            b_val = int(val[4:6], 16)
                            if r_val > 200 and g_val > 150 and b_val < 80:
                                gold_border_found = True
                                break
            except:
                pass
        if gold_border_found:
            print(f"PASS: C7 — Slide 8 has gold border shape (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: C7 — No gold border found on slide 8")
    except Exception as e:
        print(f"ERROR: C7 — {e}")

    # ---------------------------------------------------------------
    # Component 8: Blue (#1565C0) accent used in presentation (0.10 pts)
    # Initial has no blue accents; golden uses 1565C0 for title and decorative elements.
    # ---------------------------------------------------------------
    try:
        blue_found = False
        # Check slide 1 title and decorative rect
        slide1 = slides[0]
        for shape in slide1.shapes:
            # Check text color
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color.type is not None:
                                rgb = str(run.font.color.rgb).upper()
                                if rgb == '1565C0':
                                    blue_found = True
                        except:
                            pass
            # Check shape fill
            try:
                fill = shape.fill
                if fill.type == 1:
                    rgb = str(fill.fore_color.rgb).upper()
                    if rgb == '1565C0':
                        blue_found = True
            except:
                pass

        if blue_found:
            print(f"PASS: C8 — Blue #1565C0 accent found (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: C8 — Blue #1565C0 accent not found")
    except Exception as e:
        print(f"ERROR: C8 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
