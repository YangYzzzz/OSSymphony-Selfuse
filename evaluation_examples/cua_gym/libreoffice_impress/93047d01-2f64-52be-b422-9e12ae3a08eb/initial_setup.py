"""
Initial Setup: Insert a 3D pie chart on slide 7 showing market share
Task ID: impress_tm_087
Domain: libreoffice_impress

Creates a 9-slide market analysis presentation. Slide 7 has only a title 'Market Share'.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_087'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    """Add a slide with title and bullet content."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
    return slide


def add_title_only_slide(prs, title_text):
    """Add a slide with only a title (layout index 5=Blank or 6=TitleOnly)."""
    # Use blank layout and add a manual title textbox
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8.5), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    add_title_slide(prs, "Market Analysis Report", "Global Technology Sector — Q4 2025")

    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "Global tech market reached $5.3 trillion in 2025",
        "Cloud computing segment grew 22% year-over-year",
        "AI and machine learning investment surged by 45%",
        "Mobile device market stabilized at 1.4 billion units",
        "Enterprise software revenue exceeded $680 billion",
    ])

    # Slide 3: Industry Trends
    add_content_slide(prs, "Industry Trends", [
        "Generative AI integration across all software categories",
        "Edge computing adoption accelerating in manufacturing",
        "Cybersecurity spending increased to $215 billion globally",
        "Quantum computing R&D investments doubled since 2023",
        "Sustainability-driven tech solutions gaining traction",
    ])

    # Slide 4: Regional Analysis
    add_content_slide(prs, "Regional Analysis", [
        "North America: 38% market share, led by cloud & AI",
        "Asia-Pacific: 31% market share, fastest growth at 18% CAGR",
        "Europe: 22% market share, strong in automotive tech",
        "Rest of World: 9% market share, emerging fintech hubs",
    ])

    # Slide 5: Revenue Breakdown
    add_content_slide(prs, "Revenue Breakdown by Segment", [
        "Cloud Services: $1.2T (23% of total)",
        "Enterprise Software: $680B (13% of total)",
        "Consumer Electronics: $820B (15% of total)",
        "Semiconductors: $620B (12% of total)",
        "IT Services & Consulting: $1.1T (21% of total)",
        "Other Segments: $880B (16% of total)",
    ])

    # Slide 6: Competitive Landscape
    add_content_slide(prs, "Competitive Landscape", [
        "Top 5 players control 54% of cloud infrastructure",
        "Startup funding in AI reached $92B in 2025",
        "M&A activity increased 30% in cybersecurity sector",
        "Open-source adoption reshaping enterprise tool market",
    ])

    # Slide 7: Market Share — TITLE ONLY, no chart
    add_title_only_slide(prs, "Market Share")

    # Slide 8: Growth Projections
    add_content_slide(prs, "Growth Projections 2026-2030", [
        "Global tech market expected to reach $7.8T by 2030",
        "AI market projected at $1.5T by 2030 (35% CAGR)",
        "Cloud adoption to reach 85% of enterprises by 2028",
        "5G-enabled IoT devices to exceed 25 billion by 2030",
        "Green tech investments projected at $400B annually",
    ])

    # Slide 9: Recommendations
    add_content_slide(prs, "Strategic Recommendations", [
        "Prioritize AI integration across product portfolios",
        "Invest in multi-cloud and hybrid infrastructure",
        "Strengthen cybersecurity posture and compliance",
        "Expand into Asia-Pacific markets for growth",
        "Develop sustainability-focused technology solutions",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
