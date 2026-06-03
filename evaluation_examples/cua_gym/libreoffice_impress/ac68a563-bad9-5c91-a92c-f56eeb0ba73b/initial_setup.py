"""
Initial Setup: Quarterly Report Presentation with table in middle of slide 2
Task ID: osworld_impress_table_position_bottom_002
Domain: libreoffice_impress

Creates a 6-slide quarterly report deck where slide 2 has a title and
a 3x4 summary table positioned in the MIDDLE of the slide.
The agent task is to move the table to the bottom of slide 2.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_table_position_bottom_002'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_text(slide, title_text, font_size=28, bold=True):
    """Helper to set slide title text."""
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        tf.text = title_text
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(font_size)
                run.font.bold = bold
                run.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)


def create_initial():
    prs = Presentation()
    # Standard widescreen slide size: 10 x 7.5 inches
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_layouts = prs.slide_layouts
    # Layout 0 = Title Slide, 1 = Title+Content, 5 = Blank, 6 = Title Only

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(slide_layouts[0])
    slide1.shapes.title.text = "Q1 2025 Quarterly Business Review"
    slide1.placeholders[1].text = "Prepared by: Strategic Planning Division\nDate: March 2025"
    for para in slide1.placeholders[1].text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x44, 0x47, 0x2C)

    # ---- Slide 2: Executive Summary (table in MIDDLE of slide) ----
    slide2 = prs.slides.add_slide(slide_layouts[6])  # Title Only layout
    add_title_text(slide2, "Executive Summary")

    # Place a 3-row x 4-column table in the MIDDLE of the slide
    # Middle area: top ~2.2 inches, height ~2.5 inches
    rows, cols = 3, 4
    table_left   = Inches(1.0)
    table_top    = Inches(2.2)   # Middle of slide (slide height = 7.5 in)
    table_width  = Inches(8.0)
    table_height = Inches(2.5)

    tbl_shape = slide2.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
    tbl = tbl_shape.table

    # Set column widths
    tbl.columns[0].width = Inches(2.5)
    tbl.columns[1].width = Inches(2.0)
    tbl.columns[2].width = Inches(2.0)
    tbl.columns[3].width = Inches(1.5)

    # Header row (row 0)
    headers = ["Metric", "Q4 2024", "Q1 2025", "Change"]
    for c, header in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = header
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.size = Pt(14)
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # Header cell fill
        from pptx.oxml.ns import qn
        from lxml import etree
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', '1F3964')

    # Data row 1
    row1_data = ["Revenue ($M)", "128.4", "142.7", "+11.1%"]
    for c, val in enumerate(row1_data):
        cell = tbl.cell(1, c)
        cell.text = val
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(13)

    # Data row 2
    row2_data = ["Net Profit ($M)", "24.6", "29.3", "+19.1%"]
    for c, val in enumerate(row2_data):
        cell = tbl.cell(2, c)
        cell.text = val
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(13)

    # ---- Slide 3: Revenue Performance ----
    slide3 = prs.slides.add_slide(slide_layouts[1])  # Title+Content
    add_title_text(slide3, "Revenue Performance")
    content_ph = slide3.placeholders[1]
    tf = content_ph.text_frame
    tf.text = "Total Revenue: $142.7M (+11.1% YoY)"
    bullets = [
        "North America: $78.4M — up 9.3% from Q4 2024",
        "EMEA Region: $35.2M — up 15.7% from Q4 2024",
        "Asia-Pacific: $29.1M — up 8.9% from Q4 2024",
        "Key driver: New enterprise contracts in cloud services segment",
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.level = 1
        for run in p.runs:
            run.font.size = Pt(16)

    # ---- Slide 4: Operational Highlights ----
    slide4 = prs.slides.add_slide(slide_layouts[1])
    add_title_text(slide4, "Operational Highlights")
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Key milestones achieved this quarter:"
    highlights = [
        "Launched Project Atlas — cloud migration platform for SMBs",
        "Signed 47 new enterprise contracts worth $18.6M ARR",
        "Headcount grew to 2,340 full-time employees (+6.8%)",
        "Customer satisfaction (NPS) improved from 52 to 61",
        "Reduced average support ticket resolution time by 22%",
    ]
    for h in highlights:
        p = tf4.add_paragraph()
        p.text = h
        p.level = 1
        for run in p.runs:
            run.font.size = Pt(15)

    # ---- Slide 5: Financial Outlook ----
    slide5 = prs.slides.add_slide(slide_layouts[1])
    add_title_text(slide5, "Q2 2025 Financial Outlook")
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Guidance for Q2 2025:"
    outlook = [
        "Revenue forecast: $148M – $155M",
        "EBITDA margin target: 22%–24%",
        "R&D investment: $12M (focused on AI-assisted analytics)",
        "Planned headcount additions: ~120 in Engineering and Sales",
        "Key risk: Supply chain disruptions in Asia-Pacific region",
    ]
    for o in outlook:
        p = tf5.add_paragraph()
        p.text = o
        p.level = 1
        for run in p.runs:
            run.font.size = Pt(15)

    # ---- Slide 6: Closing ----
    slide6 = prs.slides.add_slide(slide_layouts[0])
    slide6.shapes.title.text = "Thank You"
    slide6.placeholders[1].text = (
        "Questions & Discussion\n\n"
        "Contact: strategy@companyexample.com\n"
        "Investor Relations: ir@companyexample.com"
    )

    # Save
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
