"""
Reward Script: Reposition table on slide 2 to near the bottom of the slide
Task ID: osworld_impress_table_position_bottom_002
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.6): Table top position is in the bottom half of the slide (top > 50% slide height)
  - Component 2 (0.4): Table top position is clearly near the bottom (top > 60% slide height)
"""

import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_table_position_bottom_002'


def verify_task(file_path):
    """
    Verify that the table on slide 2 has been moved to near the bottom of the slide.
    Initial state: table top at 2.2 in (middle of slide, slide height = 7.5 in)
    Golden state: table top at 5.0 in (clearly in bottom area)

    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load presentation {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Basic sanity: the file should have at least 2 slides
    if len(prs.slides) < 2:
        print(f"CRITICAL: Expected at least 2 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide_height = prs.slide_height  # in EMU; slide is 7.5 in = 6858000 EMU

    # Find the table on slide 2 (index 1)
    slide2 = prs.slides[1]
    table_shape = None
    for shape in slide2.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_shape = shape
            break

    if table_shape is None:
        print("FAIL: No table found on slide 2")
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    table_top = table_shape.top
    table_top_inches = table_top / 914400
    slide_height_inches = slide_height / 914400
    midpoint_emu = slide_height / 2
    sixty_pct_emu = slide_height * 0.60

    print(f"INFO: Slide height = {slide_height_inches:.3f} in ({slide_height} EMU)")
    print(f"INFO: Table top = {table_top_inches:.3f} in ({table_top} EMU)")
    print(f"INFO: Slide midpoint = {slide_height_inches/2:.3f} in ({midpoint_emu:.0f} EMU)")
    print(f"INFO: 60% threshold = {slide_height_inches*0.60:.3f} in ({sixty_pct_emu:.0f} EMU)")

    # Component 1: Table top is in the bottom half of the slide (top > 50% of slide height)
    # Initial: top = 2.2 in => FAILS (2.2 < 3.75)
    # Golden: top = 5.0 in => PASSES (5.0 > 3.75)
    try:
        if table_top > midpoint_emu:
            print(f"PASS: Component 1 — Table is in the bottom half of the slide "
                  f"(top={table_top_inches:.3f} in > midpoint={slide_height_inches/2:.3f} in) (0.6 pts)")
            total_score += 0.6
        else:
            print(f"FAIL: Component 1 — Table is NOT in the bottom half of the slide "
                  f"(top={table_top_inches:.3f} in <= midpoint={slide_height_inches/2:.3f} in)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Table top is clearly near the bottom (top > 60% of slide height)
    # Initial: top = 2.2 in => FAILS (2.2 < 4.5)
    # Golden: top = 5.0 in => PASSES (5.0 > 4.5)
    try:
        if table_top > sixty_pct_emu:
            print(f"PASS: Component 2 — Table is clearly near the bottom (>60% down) "
                  f"(top={table_top_inches:.3f} in > 60% threshold={slide_height_inches*0.60:.3f} in) (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 2 — Table is NOT clearly near the bottom (>60% down) "
                  f"(top={table_top_inches:.3f} in <= 60% threshold={slide_height_inches*0.60:.3f} in)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the VM environment
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
