"""
Reward Script: Presentation cleanup — bold+navy title on slide 2, gray background on slide 3, delete author textbox on slide 4
Task ID: osworld_impress_multi_op_combined_008
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 2 title run is bold                      (0.20 pts)
  Component 2: Slide 2 title run color is dark navy #1B3A6B   (0.20 pts)
  Component 3: Slide 3 background is gray (#D9D9D9)           (0.30 pts)
  Component 4: Slide 4 has no shape containing 'Author: Jane Smith' (0.30 pts)
  Total: 1.0
"""

import os
from pptx import Presentation
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_multi_op_combined_008'

EXPECTED_NAVY = '1B3A6B'
EXPECTED_GRAY = 'D9D9D9'
AUTHOR_TEXT = 'Author: Jane Smith'


def persist_app_state():
    """Best-effort save hook: send Ctrl+S to any open LibreOffice window."""
    try:
        import pyautogui
        import time
        os.environ["DISPLAY"] = ":0"
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent to LibreOffice Impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Sanity gate: must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"CRITICAL: Expected at least 4 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # -----------------------------------------------------------------
    # Component 1: Slide 2 title run is bold (0.20 pts)
    # In initial_env: Bold=False; in golden_env: Bold=True
    # -----------------------------------------------------------------
    try:
        slide2 = prs.slides[1]
        title_shape = None
        for shape in slide2.shapes:
            if shape.has_text_frame and hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                if shape.placeholder_format.idx == 0:
                    title_shape = shape
                    break
        if title_shape is None:
            # Fallback: use the first shape with a text frame named like 'Title'
            for shape in slide2.shapes:
                if shape.has_text_frame and 'title' in shape.name.lower():
                    title_shape = shape
                    break
        if title_shape is None:
            print("FAIL: Component 1 — could not locate title shape on slide 2")
        else:
            title_bold = False
            for para in title_shape.text_frame.paragraphs:
                nonempty = [r for r in para.runs if (r.text or "").strip()]
                for run in nonempty:
                    if run.font.bold is True:
                        title_bold = True
                        break
                if title_bold:
                    break
            if title_bold:
                print("PASS: Component 1 — slide 2 title is bold (0.20 pts)")
                total_score += 0.20
            else:
                print("FAIL: Component 1 — slide 2 title is NOT bold")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # -----------------------------------------------------------------
    # Component 2: Slide 2 title run color is dark navy #1B3A6B (0.20 pts)
    # In initial_env: Color=000000; in golden_env: Color=1B3A6B
    # -----------------------------------------------------------------
    try:
        slide2 = prs.slides[1]
        title_shape = None
        for shape in slide2.shapes:
            if shape.has_text_frame and hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                if shape.placeholder_format.idx == 0:
                    title_shape = shape
                    break
        if title_shape is None:
            for shape in slide2.shapes:
                if shape.has_text_frame and 'title' in shape.name.lower():
                    title_shape = shape
                    break
        if title_shape is None:
            print("FAIL: Component 2 — could not locate title shape on slide 2")
        else:
            navy_found = False
            for para in title_shape.text_frame.paragraphs:
                nonempty = [r for r in para.runs if (r.text or "").strip()]
                for run in nonempty:
                    if run.font.color.type is not None:
                        color_str = str(run.font.color.rgb).upper()
                        if color_str == EXPECTED_NAVY:
                            navy_found = True
                            break
                if navy_found:
                    break
            if navy_found:
                print(f"PASS: Component 2 — slide 2 title color is #{EXPECTED_NAVY} (0.20 pts)")
                total_score += 0.20
            else:
                # Show what color was actually set
                found_colors = []
                for para in title_shape.text_frame.paragraphs:
                    for run in para.runs:
                        if (run.text or "").strip():
                            if run.font.color.type is not None:
                                found_colors.append(str(run.font.color.rgb).upper())
                            else:
                                found_colors.append("(inherited)")
                print(f"FAIL: Component 2 — slide 2 title color is not #{EXPECTED_NAVY}, found: {found_colors}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # -----------------------------------------------------------------
    # Component 3: Slide 3 background is gray (#D9D9D9) (0.30 pts)
    # In initial_env: FFFFFF (white); in golden_env: D9D9D9 (gray)
    # -----------------------------------------------------------------
    try:
        slide3 = prs.slides[2]
        fill = slide3.background.fill
        bg_color = None
        if fill.type == 1:  # SOLID
            bg_color = str(fill.fore_color.rgb).upper()
        elif fill.type == 5:  # inherited from master
            master_fill = slide3.slide_layout.slide_master.background.fill
            if master_fill.type == 1:
                bg_color = str(master_fill.fore_color.rgb).upper()

        if bg_color == EXPECTED_GRAY:
            print(f"PASS: Component 3 — slide 3 background is gray #{EXPECTED_GRAY} (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 3 — slide 3 background is not gray #{EXPECTED_GRAY}, found: {bg_color}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # -----------------------------------------------------------------
    # Component 4: Slide 4 has no shape containing 'Author: Jane Smith' (0.30 pts)
    # In initial_env: TextBox 3 with 'Author: Jane Smith' exists
    # In golden_env: that text box is deleted
    # -----------------------------------------------------------------
    try:
        slide4 = prs.slides[3]
        author_found = False
        for shape in slide4.shapes:
            if shape.has_text_frame:
                # Collect all text from this shape (including grouped shapes)
                full_text = ""
                for para in shape.text_frame.paragraphs:
                    full_text += para.text
                if AUTHOR_TEXT in full_text:
                    author_found = True
                    print(f"FAIL: Component 4 — shape '{shape.name}' still contains '{AUTHOR_TEXT}'")
                    break

        if not author_found:
            print(f"PASS: Component 4 — no shape with '{AUTHOR_TEXT}' found on slide 4 (0.30 pts)")
            total_score += 0.30
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # -----------------------------------------------------------------
    final_score = min(round(total_score, 4), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point: resolve canonical file path and run verification
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
