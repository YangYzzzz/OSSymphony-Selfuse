"""
Initial Setup: Research report presentation with 6 slides
Task ID: osworld_impress_multi_op_combined_008
Domain: libreoffice_impress

Creates the initial state: 6-slide research report deck where:
  - Slide 2 title is plain black (not bold, not navy)
  - Slide 3 has a white/no background
  - Slide 4 has a text box with 'Author: Jane Smith'
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_multi_op_combined_008'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard widescreen layout
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Climate Change Impact Assessment"
    slide1.placeholders[1].text = "A Comprehensive Research Report\nEnvironmental Studies Division, 2025"

    # --- Slide 2: Introduction (title must be plain black, NOT bold, NOT navy) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    title2.text = "Introduction & Background"
    # Explicitly set title to plain black, not bold
    for para in title2.text_frame.paragraphs:
        for run in para.runs:
            run.font.bold = False
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Plain black

    content2 = slide2.placeholders[1]
    content2.text = (
        "Global temperatures have risen by 1.1°C since pre-industrial times.\n"
        "Extreme weather events have increased in frequency and severity.\n"
        "Arctic sea ice coverage has decreased by 13% per decade.\n"
        "Rising sea levels threaten coastal communities worldwide.\n"
        "Ecosystem disruption affects biodiversity across all continents."
    )

    # --- Slide 3: Methodology (white background — achieved by no fill set) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = "Research Methodology"
    # Explicitly set white background on slide 3
    fill3 = slide3.background.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White background

    content3 = slide3.placeholders[1]
    content3.text = (
        "Data Collection: 47 monitoring stations across 6 continents.\n"
        "Time Period: January 2015 – December 2024 (10-year span).\n"
        "Statistical Analysis: Regression models and time-series analysis.\n"
        "Peer Review: Validated by 12 independent research institutions.\n"
        "Satellite Data: Integrated NOAA and ESA remote sensing datasets."
    )

    # --- Slide 4: Key Findings (has text box with 'Author: Jane Smith') ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    title4.text = "Key Findings"

    content4 = slide4.placeholders[1]
    content4.text = (
        "Temperature anomalies exceed 2.5°C in polar regions.\n"
        "Ocean acidification has increased by 26% since 1850.\n"
        "Coral reef degradation affects 50% of monitored reefs.\n"
        "Permafrost thaw releasing 1.5 Gt CO₂ equivalent annually.\n"
        "Species extinction rate is 100x the natural background rate."
    )

    # Add author textbox — this must be removed in the golden patch
    txBox = slide4.shapes.add_textbox(
        Inches(0.5), Inches(6.6), Inches(4.0), Inches(0.6)
    )
    tf = txBox.text_frame
    tf.text = "Author: Jane Smith"
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(12)
            run.font.italic = True
            run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # --- Slide 5: Data Analysis ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    title5.text = "Data Analysis & Trends"

    content5 = slide5.placeholders[1]
    content5.text = (
        "Linear regression shows R² = 0.94 for temperature increase trend.\n"
        "Precipitation patterns show 18% variance from 20-year baseline.\n"
        "Drought frequency in Mediterranean region increased by 34%.\n"
        "Storm surge events in coastal regions up by 41% since 2000.\n"
        "Carbon sequestration capacity of forests reduced by 22%."
    )

    # --- Slide 6: Conclusions & Recommendations ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    title6.text = "Conclusions & Recommendations"

    content6 = slide6.placeholders[1]
    content6.text = (
        "Immediate reduction of greenhouse gas emissions is critical.\n"
        "Investment in renewable energy must increase by 300% by 2035.\n"
        "Coastal infrastructure requires urgent adaptation measures.\n"
        "International cooperation frameworks must be strengthened.\n"
        "Continued monitoring and research funding is essential."
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
