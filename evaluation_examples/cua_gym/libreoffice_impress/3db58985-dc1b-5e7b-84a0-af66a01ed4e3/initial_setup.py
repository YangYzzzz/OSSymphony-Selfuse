"""
Initial Setup: Create an 8-slide board meeting presentation with no transitions.
Task ID: impress_tm_042
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_042'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Q1 2025 Board Meeting"
    slide.placeholders[1].text = "Pinnacle Financial Holdings\nMarch 28, 2025"
    return slide


def add_content_slide(prs, title, bullet_points):
    """Add a content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title ---
    add_title_slide(prs)

    # --- Slide 2: Executive Summary ---
    add_content_slide(prs, "Executive Summary", [
        "Total revenue reached $142.3M, up 18% year-over-year",
        "Net income improved to $31.7M with 22.3% margin",
        "Customer acquisition cost reduced by 12% to $47.20",
        "Employee headcount grew to 1,284 across 6 offices",
        "Three new enterprise partnerships signed in Q1",
    ])

    # --- Slide 3: Financial Overview ---
    add_content_slide(prs, "Financial Overview", [
        "Gross revenue: $142.3M (Q4 2024: $128.1M)",
        "Operating expenses: $98.6M including $12.4M R&D",
        "EBITDA: $43.7M representing 30.7% margin",
        "Free cash flow: $28.9M after capital expenditures",
        "Cash reserves: $215.4M with $50M credit facility available",
    ])

    # --- Slide 4: Product & Technology ---
    add_content_slide(prs, "Product & Technology Update", [
        "Platform 3.0 launched February 12 with AI-driven analytics",
        "Mobile app downloads surpassed 2.1M lifetime installs",
        "API response time improved 34% to average 89ms",
        "Security audit completed with zero critical findings",
        "Patent filed for predictive risk assessment algorithm",
    ])

    # --- Slide 5: Market Expansion ---
    add_content_slide(prs, "Market Expansion Strategy", [
        "European market entry planned for Q3 via London office",
        "APAC partnership with Sumitomo Financial Group finalized",
        "SMB segment grew 27% with 340 new accounts",
        "Government contracts pipeline valued at $18.5M",
        "Brand awareness increased to 42% in target demographics",
    ])

    # --- Slide 6: Risk & Compliance ---
    add_content_slide(prs, "Risk & Compliance", [
        "SOC 2 Type II certification renewed through March 2026",
        "GDPR compliance framework updated for new EU guidelines",
        "Cybersecurity insurance coverage expanded to $100M",
        "Regulatory change impact assessment completed for 2025",
        "Internal audit identified 3 low-severity process gaps",
    ])

    # --- Slide 7: Talent & Culture ---
    add_content_slide(prs, "Talent & Culture", [
        "Employee satisfaction score: 4.3/5.0 (industry avg: 3.8)",
        "Voluntary turnover rate: 8.2% vs 14.1% industry benchmark",
        "Leadership development program enrolled 45 participants",
        "Diversity hiring targets exceeded: 52% of new hires",
        "Remote work policy expanded to include 4-day office weeks",
    ])

    # --- Slide 8: Closing ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[0])
    slide8.shapes.title.text = "Questions & Next Steps"
    slide8.placeholders[1].text = "Thank you for your continued partnership\nNext board meeting: June 27, 2025"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
