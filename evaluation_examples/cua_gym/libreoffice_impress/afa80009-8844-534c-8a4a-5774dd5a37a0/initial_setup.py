#!/usr/bin/env python3
"""
initial_setup.py - Creates a 20-slide Training_Course.pptx with realistic content.
No navigation buttons on any slide.
"""

import subprocess, shlex, os, time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

OUTPUT_PATH = "/home/user/impress_fix_049.pptx"

# Define 20 slides of training course content
SLIDES = [
    {
        "layout": 0,  # Title Slide
        "title": "Advanced Data Analytics Training Course",
        "subtitle": "Building Skills for Data-Driven Decision Making\nQ2 2026 | Corporate Learning & Development",
    },
    {
        "layout": 1,  # Title + Content
        "title": "Course Agenda",
        "bullets": [
            "Module 1: Introduction to Data Analytics",
            "Module 2: Data Collection and Cleaning",
            "Module 3: Exploratory Data Analysis",
            "Module 4: Statistical Methods",
            "Module 5: Data Visualization",
            "Module 6: Machine Learning Fundamentals",
            "Module 7: Practical Case Studies",
            "Module 8: Final Assessment",
        ],
    },
    {
        "layout": 1,
        "title": "Course Objectives",
        "bullets": [
            "Understand core data analytics concepts and terminology",
            "Master data cleaning and preparation techniques",
            "Apply statistical methods to real-world datasets",
            "Create effective data visualizations for stakeholders",
            "Build foundational machine learning models",
            "Develop a data-driven mindset for business decisions",
        ],
    },
    {
        "layout": 1,
        "title": "Module 1: Introduction to Data Analytics",
        "bullets": [
            "What is data analytics and why does it matter?",
            "Types of analytics: Descriptive, Diagnostic, Predictive, Prescriptive",
            "The data analytics lifecycle",
            "Key tools and technologies overview",
            "Industry applications and career paths",
        ],
    },
    {
        "layout": 1,
        "title": "Module 2: Data Collection and Cleaning",
        "bullets": [
            "Data sources: databases, APIs, web scraping, surveys",
            "Data quality dimensions: accuracy, completeness, consistency",
            "Common data cleaning techniques",
            "Handling missing values and outliers",
            "Data transformation and normalization",
            "Best practices for data documentation",
        ],
    },
    {
        "layout": 1,
        "title": "Exercise 1: Data Cleaning Workshop",
        "bullets": [
            "Download the sample customer dataset from the shared drive",
            "Identify and document data quality issues",
            "Apply cleaning techniques to resolve issues",
            "Validate your cleaned dataset against the reference",
            "Time allotted: 45 minutes",
            "Work in pairs and discuss your approach",
        ],
    },
    {
        "layout": 1,
        "title": "Module 3: Exploratory Data Analysis",
        "bullets": [
            "Summary statistics and distributions",
            "Correlation analysis and feature relationships",
            "Pattern recognition and anomaly detection",
            "EDA workflow and best practices",
            "Tools: Python pandas, R, Excel pivot tables",
        ],
    },
    {
        "layout": 1,
        "title": "Module 4: Statistical Methods",
        "bullets": [
            "Descriptive statistics: mean, median, mode, standard deviation",
            "Probability distributions and sampling",
            "Hypothesis testing and confidence intervals",
            "Regression analysis: linear and logistic",
            "ANOVA and chi-square tests",
            "Common pitfalls in statistical analysis",
        ],
    },
    {
        "layout": 1,
        "title": "Exercise 2: Statistical Analysis Lab",
        "bullets": [
            "Analyze the provided sales dataset",
            "Calculate descriptive statistics for key metrics",
            "Perform hypothesis testing on regional differences",
            "Build a simple regression model",
            "Present findings to your team",
            "Time allotted: 60 minutes",
        ],
    },
    {
        "layout": 1,
        "title": "Module 5: Data Visualization",
        "bullets": [
            "Principles of effective data visualization",
            "Choosing the right chart type for your data",
            "Color theory and accessibility considerations",
            "Dashboard design principles",
            "Tools: Matplotlib, Seaborn, Tableau, Power BI",
        ],
    },
    {
        "layout": 1,
        "title": "Visualization Best Practices",
        "bullets": [
            "Keep it simple: avoid chartjunk and unnecessary decoration",
            "Use appropriate scales and axes labels",
            "Tell a story with your data",
            "Consider your audience's technical level",
            "Always include context and annotations",
            "Test for colorblind accessibility",
        ],
    },
    {
        "layout": 1,
        "title": "Module 6: Machine Learning Fundamentals",
        "bullets": [
            "Supervised vs. unsupervised learning",
            "Classification and regression problems",
            "Model training, validation, and testing",
            "Feature engineering and selection",
            "Model evaluation metrics",
            "Overfitting and regularization",
        ],
    },
    {
        "layout": 1,
        "title": "Common ML Algorithms Overview",
        "bullets": [
            "Decision Trees and Random Forests",
            "Support Vector Machines (SVM)",
            "K-Nearest Neighbors (KNN)",
            "K-Means Clustering",
            "Neural Networks basics",
            "When to use which algorithm",
        ],
    },
    {
        "layout": 1,
        "title": "Exercise 3: Build Your First ML Model",
        "bullets": [
            "Use the provided customer churn dataset",
            "Perform feature engineering and selection",
            "Train a classification model using scikit-learn",
            "Evaluate model performance with appropriate metrics",
            "Compare at least two different algorithms",
            "Time allotted: 90 minutes",
        ],
    },
    {
        "layout": 1,
        "title": "Module 7: Practical Case Studies",
        "bullets": [
            "Case Study 1: Customer segmentation for marketing",
            "Case Study 2: Sales forecasting for inventory management",
            "Case Study 3: Fraud detection in financial transactions",
            "Case Study 4: Sentiment analysis of product reviews",
            "Lessons learned and common challenges",
        ],
    },
    {
        "layout": 1,
        "title": "Case Study Deep Dive: Customer Segmentation",
        "bullets": [
            "Business objective: Targeted marketing campaigns",
            "Data sources: Transaction history, demographics, behavior",
            "Approach: RFM analysis and K-Means clustering",
            "Results: 5 distinct customer segments identified",
            "Impact: 23% increase in campaign conversion rate",
            "Key takeaways for similar projects",
        ],
    },
    {
        "layout": 1,
        "title": "Module 8: Final Assessment",
        "bullets": [
            "Individual project: End-to-end data analytics pipeline",
            "Choose from provided datasets or propose your own",
            "Requirements: Data cleaning, EDA, modeling, visualization",
            "20-minute presentation to the class",
            "Peer review and feedback session",
            "Grading rubric available on the course portal",
        ],
    },
    {
        "layout": 1,
        "title": "Course Resources and References",
        "bullets": [
            "Course materials: Available on the Learning Management System",
            "Recommended books: 'Python for Data Analysis' by Wes McKinney",
            "Online resources: Kaggle, DataCamp, Coursera",
            "Office hours: Tuesdays and Thursdays, 3-5 PM",
            "Discussion forum: Use the #data-analytics Slack channel",
            "Additional practice datasets on the shared drive",
        ],
    },
    {
        "layout": 1,
        "title": "Key Takeaways",
        "bullets": [
            "Data quality is the foundation of good analytics",
            "Always start with exploratory analysis before modeling",
            "Choose the simplest model that solves the problem",
            "Visualization is crucial for communicating insights",
            "Continuous learning is essential in this rapidly evolving field",
            "Practice with real datasets to build confidence",
        ],
    },
    {
        "layout": 1,
        "title": "Thank You & Q&A",
        "bullets": [
            "Thank you for your participation and engagement!",
            "Questions and open discussion",
            "Contact: training@company.com",
            "Feedback survey: Please complete before leaving",
            "Certificates will be issued upon course completion",
            "Next advanced course starts in September 2026",
        ],
    },
]


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, slide_data in enumerate(SLIDES):
        layout_idx = slide_data["layout"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data["title"]

        if layout_idx == 0:
            # Title slide - set subtitle
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_data["subtitle"]
        elif layout_idx == 1:
            # Content slide - set bullets
            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for j, bullet in enumerate(slide_data["bullets"]):
                    if j == 0:
                        tf.paragraphs[0].text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 0

    prs.save(OUTPUT_PATH)
    print(f"Saved presentation to {OUTPUT_PATH}")


def launch_gui(command, delay_sec=2.0):
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


if __name__ == "__main__":
    build_presentation()
    launch_gui(f'libreoffice --impress "{OUTPUT_PATH}"', delay_sec=2.0)
    print("LibreOffice Impress launched.")
