"""
Reward Script: KPI Dashboard Slide for Annual_Review_2025.pptx
Task ID: impress_ps_017
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Four metric card background shapes exist on slide 3
  Component 2 (0.30): All four metric values present with correct text and bold formatting
  Component 3 (0.20): All four metric labels present (Revenue, Customers, NPS, Churn)
  Component 4 (0.20): All four change indicators present with green color and arrow
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_017'

# Expected metric data
EXPECTED_METRICS = [
    {"value": "$12.5M", "label": "Revenue", "change": "+18%"},
    {"value": "2,450", "label": "Customers", "change": "+32%"},
    {"value": "72", "label": "NPS", "change": "+8"},
    {"value": "3.2%", "label": "Churn", "change": "-1.5%"},
]


def get_all_text_shapes(slide):
    """Get all text-bearing shapes from a slide, including inside groups."""
    results = []
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                if hasattr(sub, 'text_frame'):
                    results.append(sub)
    return results


def get_shape_full_text(shape):
    """Get concatenated text from all paragraphs/runs."""
    if not hasattr(shape, 'text_frame'):
        return ""
    return shape.text_frame.text.strip()


def check_value_bold(shape, val):
    """Check if val text exists in shape with bold formatting. Returns 1 if bold, 0.5 if present but not bold, 0 if absent."""
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if val in (run.text or "") and run.font.bold:
                return 1
    return 0.5  # present but bold not confirmed


def check_green_arrow(shape):
    """Check if shape contains a green-colored arrow indicator. Returns (has_green, has_arrow)."""
    green_found = False
    arrow_found = False
    arrow_chars = ['↑', '↓', '▲', '▼', '\u2191', '\u2193']
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run_text = run.text or ""
            if any(c in run_text for c in arrow_chars):
                arrow_found = arrow_found or len(run_text.strip()) > 0
            try:
                if run.font.color.type is not None:
                    rgb = str(run.font.color.rgb).upper()
                    r_val = int(rgb[0:2], 16)
                    g_val = int(rgb[2:4], 16)
                    b_val = int(rgb[4:6], 16)
                    if g_val > r_val and g_val > b_val:
                        green_found = green_found or (g_val > 0)
            except Exception:
                pass
    return green_found, arrow_found


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Pt, Emu
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file has at least 3 slides
    if len(prs.slides) < 3:
        print(f"PRECONDITION FAIL: Need at least 3 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[2]  # Slide 3 (0-indexed)
    text_shapes = get_all_text_shapes(slide)

    # Pre-existing shapes on slide 3 (should not count as task changes)
    PRE_EXISTING_TEXTS = {"2025 Key Metrics", "Key performance indicators for fiscal year 2025"}

    # Get non-pre-existing text shapes (the KPI cards added by the task)
    new_text_shapes = [s for s in text_shapes if get_shape_full_text(s) not in PRE_EXISTING_TEXTS]

    # Count auto shapes (card backgrounds) on slide 3
    auto_shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]

    # Component 1: Four card background shapes exist on slide 3 (0.30 points)
    # Initial state has 0 auto shapes on slide 3; golden has 4
    try:
        num_auto = len(auto_shapes)
        if num_auto >= 4:
            print(f"PASS: Component 1 - Found {num_auto} card background shapes (>= 4) (0.30 pts)")
            total_score += 0.30
        elif num_auto >= 2:
            print(f"PARTIAL: Component 1 - Found {num_auto} card backgrounds (need 4) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 - Found {num_auto} card background shapes, need 4")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: All four metric values present with bold formatting (0.30 points)
    # Expected values: $12.5M, 2,450, 72, 3.2%
    try:
        found_values = 0
        for metric in EXPECTED_METRICS:
            val = metric["value"]
            matched = False
            for shape in new_text_shapes:
                if val in get_shape_full_text(shape):
                    credit = check_value_bold(shape, val)
                    found_values += credit
                    if credit >= 1:
                        print(f"  Found value '{val}' with bold formatting")
                    else:
                        print(f"  Found value '{val}' but bold formatting not confirmed")
                    matched = credit > 0  # mark as found
                    break
            if not matched:
                print(f"  Missing value '{val}'")

        if found_values >= 4:
            print(f"PASS: Component 2 - All 4 metric values found with bold (0.30 pts)")
            total_score += 0.30
        elif found_values > 0:
            value_score = round((found_values / 4.0) * 0.30, 4)
            print(f"PARTIAL: Component 2 - {found_values}/4 metric values found ({value_score} pts)")
            total_score += value_score
        else:
            print(f"FAIL: Component 2 - No metric values found")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: All four metric labels present (0.20 points)
    # Expected labels: Revenue, Customers, NPS, Churn
    try:
        found_labels = 0
        for metric in EXPECTED_METRICS:
            label = metric["label"]
            for shape in new_text_shapes:
                if label.lower() in get_shape_full_text(shape).lower():
                    found_labels += 1
                    print(f"  Found label '{label}'")
                    break

        if found_labels >= 4:
            print(f"PASS: Component 3 - All 4 metric labels found (0.20 pts)")
            total_score += 0.20
        elif found_labels > 0:
            label_score = round((found_labels / 4.0) * 0.20, 4)
            print(f"PARTIAL: Component 3 - {found_labels}/4 metric labels found ({label_score} pts)")
            total_score += label_score
        else:
            print(f"FAIL: Component 3 - No metric labels found")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Change indicators with green color and arrow (0.20 points)
    # Expected: +18%, +32%, +8, -1.5% all with green color and arrow symbol
    try:
        found_changes = 0
        for metric in EXPECTED_METRICS:
            change = metric["change"]
            for shape in new_text_shapes:
                if change in get_shape_full_text(shape):
                    has_green, has_arrow = check_green_arrow(shape)
                    if has_green and has_arrow:
                        found_changes += 1
                        print(f"  Found change '{change}' with green arrow")
                    elif has_green or has_arrow:
                        found_changes += 0.5
                        print(f"  Found change '{change}' (green={has_green}, arrow={has_arrow})")
                    else:
                        print(f"  Found change '{change}' text but no green color or arrow")
                    break

        if found_changes >= 4:
            print(f"PASS: Component 4 - All 4 change indicators with green arrows (0.20 pts)")
            total_score += 0.20
        elif found_changes > 0:
            change_score = round((found_changes / 4.0) * 0.20, 4)
            print(f"PARTIAL: Component 4 - {found_changes}/4 change indicators ({change_score} pts)")
            total_score += change_score
        else:
            print(f"FAIL: Component 4 - No change indicators found")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
