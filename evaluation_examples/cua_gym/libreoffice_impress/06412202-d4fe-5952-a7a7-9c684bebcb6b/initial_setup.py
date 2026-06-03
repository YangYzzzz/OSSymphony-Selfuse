"""
Initial Setup: Create Annual_Review_2025.pptx with 12 slides, slide 3 has title only
Task ID: impress_ps_017
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_017'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_text(slide, text, font_size=32, bold=True, color=None):
    """Set the title placeholder text with formatting."""
    title = slide.shapes.title
    title.text = text
    for run in title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color


def add_body_textbox(slide, left, top, width, height, lines, font_size=18, color=None):
    """Add a textbox with multiple lines."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.space_after = Pt(6)
        for run in p.runs:
            run.font.size = Pt(font_size)
            if color:
                run.font.color.rgb = color


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color palette
    dark_blue = RGBColor(0x1B, 0x2A, 0x4A)
    accent_blue = RGBColor(0x2E, 0x75, 0xB6)
    white = RGBColor(0xFF, 0xFF, 0xFF)
    dark_gray = RGBColor(0x33, 0x33, 0x33)
    light_gray = RGBColor(0x66, 0x66, 0x66)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = dark_blue
    txBox = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Annual Review 2025"
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = white
    p2 = tf.add_paragraph()
    p2.text = "Meridian Technologies Inc."
    p2.alignment = PP_ALIGN.CENTER
    for run in p2.runs:
        run.font.size = Pt(24)
        run.font.color.rgb = RGBColor(0xAA, 0xCC, 0xEE)
    p3 = tf.add_paragraph()
    p3.text = "Board of Directors Presentation | March 2025"
    p3.alignment = PP_ALIGN.CENTER
    for run in p3.runs:
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x88, 0xAA, 0xCC)

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    txTitle = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = txTitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Agenda"
    for run in p.runs:
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = dark_blue
    add_body_textbox(slide2, Inches(1.2), Inches(1.8), Inches(10), Inches(5), [
        "1. Executive Summary",
        "2. 2025 Key Metrics",
        "3. Revenue Breakdown by Region",
        "4. Customer Growth Analysis",
        "5. Product Development Highlights",
        "6. Marketing Campaign Results",
        "7. Employee Engagement & Hiring",
        "8. Technology Infrastructure Updates",
        "9. Risk Assessment & Mitigation",
        "10. Financial Outlook for 2026",
        "11. Q&A",
    ], font_size=20, color=dark_gray)

    # --- Slide 3: 2025 Key Metrics (EMPTY - no metric cards) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    txTitle3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(1))
    tf = txTitle3.text_frame
    p = tf.paragraphs[0]
    p.text = "2025 Key Metrics"
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = dark_blue
    # Subtitle line
    txSub = slide3.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11), Inches(0.5))
    tf2 = txSub.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Key performance indicators for fiscal year 2025"
    for run in p2.runs:
        run.font.size = Pt(14)
        run.font.color.rgb = light_gray
    # NO metric cards - this is the empty initial state

    # --- Slide 4: Revenue Breakdown ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    txT4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = txT4.text_frame
    p = tf.paragraphs[0]
    p.text = "Revenue Breakdown by Region"
    for run in p.runs:
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = dark_blue
    # Table for revenue data
    table_shape = slide4.shapes.add_table(6, 4, Inches(1), Inches(1.8), Inches(11), Inches(3.5))
    table = table_shape.table
    headers = ["Region", "2024 Revenue", "2025 Revenue", "Growth"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    data = [
        ["North America", "$4.2M", "$5.1M", "+21%"],
        ["Europe", "$2.8M", "$3.3M", "+18%"],
        ["Asia Pacific", "$1.9M", "$2.4M", "+26%"],
        ["Latin America", "$0.8M", "$1.0M", "+25%"],
        ["Middle East & Africa", "$0.5M", "$0.7M", "+40%"],
    ]
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 5: Customer Growth ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    txT5 = slide5.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = txT5.text_frame
    p = tf.paragraphs[0]
    p.text = "Customer Growth Analysis"
    for run in p.runs:
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = dark_blue
    add_body_textbox(slide5, Inches(1), Inches(1.8), Inches(10), Inches(4.5), [
        "Total active customers: 2,450 (up from 1,856 in 2024)",
        "Enterprise accounts grew by 45% year-over-year",
        "SMB segment added 320 new customers in Q3-Q4",
        "Customer lifetime value increased to $18,400 average",
        "Referral program generated 28% of new signups",
        "Highest growth in fintech vertical (+62%)",
    ], font_size=18, color=dark_gray)

    # --- Slide 6: Product Development ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    txT6 = slide6.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = txT6.text_frame
    p = tf.paragraphs[0]
    p.text = "Product Development Highlights"
    for run in p.runs:
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = dark_blue
    add_body_textbox(slide6, Inches(1), Inches(1.8), Inches(10), Inches(4.5), [
        "Launched Meridian Analytics Pro in Q1 2025",
        "API response time improved by 40% (avg 120ms)",
        "Mobile app redesign completed - 4.8 star rating",
        "Introduced AI-powered recommendation engine",
        "99.97% platform uptime achieved across all regions",
        "12 major releases shipped on schedule",
    ], font_size=18, color=dark_gray)

    # --- Slide 7: Marketing Campaigns ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    txT7 = slide7.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = txT7.text_frame
    p = tf.paragraphs[0]
    p.text = "Marketing Campaign Results"
    for run in p.runs:
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = dark_blue
    add_body_textbox(slide7, Inches(1), Inches(1.8), Inches(10), Inches(4.5), [
        "Digital advertising spend: $1.8M (ROI: 340%)",
        "Content marketing generated 45,000 qualified leads",
        "Webinar series: 12 events, 8,200 attendees total",
        "Brand awareness increased 28% (Nielsen survey)",
        "Social media followers grew to 125,000 (+55%)",
        "Trade show presence at 6 major industry events",
    ], font_size=18, color=dark_gray)

    # --- Slide 8: Employee Engagement ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    txT8 = slide8.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = txT8.text_frame
    p = tf.paragraphs[0]
    p.text = "Employee Engagement & Hiring"
    for run in p.runs:
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = dark_blue
    add_body_textbox(slide8, Inches(1), Inches(1.8), Inches(10), Inches(4.5), [
        "Total headcount: 342 (net +68 from 2024)",
        "Employee satisfaction score: 4.3/5.0",
        "Voluntary turnover rate: 8.2% (industry avg: 13%)",
        "Engineering team expanded by 35 senior hires",
        "Launched mentorship program with 89% participation",
        "Average tenure increased to 3.4 years",
    ], font_size=18, color=dark_gray)

    # --- Slide 9: Technology Infrastructure ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    txT9 = slide9.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = txT9.text_frame
    p = tf.paragraphs[0]
    p.text = "Technology Infrastructure Updates"
    for run in p.runs:
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = dark_blue
    add_body_textbox(slide9, Inches(1), Inches(1.8), Inches(10), Inches(4.5), [
        "Completed migration to multi-cloud architecture",
        "Data processing capacity increased 3x with new GPU clusters",
        "Zero critical security incidents in 2025",
        "SOC 2 Type II certification renewed",
        "Implemented zero-trust network architecture",
        "Infrastructure cost reduced 22% through optimization",
    ], font_size=18, color=dark_gray)

    # --- Slide 10: Risk Assessment ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    txT10 = slide10.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = txT10.text_frame
    p = tf.paragraphs[0]
    p.text = "Risk Assessment & Mitigation"
    for run in p.runs:
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = dark_blue
    add_body_textbox(slide10, Inches(1), Inches(1.8), Inches(10), Inches(4.5), [
        "Market competition: Monitoring 3 emerging competitors",
        "Regulatory changes: GDPR compliance team expanded",
        "Supply chain: Diversified cloud providers to 3 vendors",
        "Talent retention: Enhanced compensation packages in Q2",
        "Cybersecurity: Quarterly penetration testing ongoing",
        "Economic downturn: 18-month operational runway maintained",
    ], font_size=18, color=dark_gray)

    # --- Slide 11: Financial Outlook ---
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    txT11 = slide11.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = txT11.text_frame
    p = tf.paragraphs[0]
    p.text = "Financial Outlook for 2026"
    for run in p.runs:
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = dark_blue
    add_body_textbox(slide11, Inches(1), Inches(1.8), Inches(10), Inches(4.5), [
        "Projected revenue: $15.8M (+26% year-over-year)",
        "Target EBITDA margin: 18% (up from 14%)",
        "Planned R&D investment: $3.2M",
        "Expected headcount growth: +45 positions",
        "New market entry: Southeast Asia expansion planned for Q2",
        "Series C fundraise target: $40M at $280M valuation",
    ], font_size=18, color=dark_gray)

    # --- Slide 12: Q&A / Thank You ---
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide12.background.fill
    fill.solid()
    fill.fore_color.rgb = dark_blue
    txBox = slide12.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = white
    p2 = tf.add_paragraph()
    p2.text = "Questions & Discussion"
    p2.alignment = PP_ALIGN.CENTER
    for run in p2.runs:
        run.font.size = Pt(24)
        run.font.color.rgb = RGBColor(0xAA, 0xCC, 0xEE)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
