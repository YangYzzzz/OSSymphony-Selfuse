"""
Initial Setup: Apply bold formatting to all title textboxes in this presentation.
Task ID: osworld_impress_title_selective_formatting_002
Domain: libreoffice_impress

Creates a 6-slide academic lecture deck where all title placeholders are in
regular (non-bold) weight. The agent task is to bold all title placeholders.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_selective_formatting_002'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # Standard widescreen dimensions (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide data: (title, body_lines)
    slides_data = [
        (
            "Introduction to Machine Learning",
            [
                "What is Machine Learning?",
                "ML enables systems to learn and improve from experience",
                "Three main types: Supervised, Unsupervised, Reinforcement",
                "Applications: image recognition, NLP, recommendation systems",
                "Course overview and learning objectives",
            ]
        ),
        (
            "Supervised Learning Fundamentals",
            [
                "Training data consists of labeled input-output pairs",
                "Algorithm learns a mapping function f(x) → y",
                "Common algorithms: Linear Regression, Decision Trees, SVM",
                "Evaluation metrics: accuracy, precision, recall, F1-score",
                "Bias-variance tradeoff is a central concern",
            ]
        ),
        (
            "Neural Networks and Deep Learning",
            [
                "Inspired by the structure of the human brain",
                "Composed of layers: input, hidden, and output",
                "Activation functions: ReLU, Sigmoid, Tanh",
                "Backpropagation algorithm for weight updates",
                "Deep networks can model highly complex functions",
            ]
        ),
        (
            "Natural Language Processing",
            [
                "Computational techniques for human language understanding",
                "Tokenization, stemming, lemmatization as preprocessing steps",
                "Word embeddings: Word2Vec, GloVe, FastText",
                "Transformer models: BERT, GPT, T5",
                "Tasks: sentiment analysis, machine translation, QA",
            ]
        ),
        (
            "Model Evaluation and Validation",
            [
                "Train/validation/test split: 70/15/15 rule of thumb",
                "Cross-validation for robust performance estimation",
                "Confusion matrix and ROC curves for classification",
                "Overfitting detection via learning curves",
                "Hyperparameter tuning with grid search and random search",
            ]
        ),
        (
            "Ethics and Fairness in AI",
            [
                "Algorithmic bias can perpetuate societal inequalities",
                "Fairness metrics: demographic parity, equalized odds",
                "Explainability: LIME, SHAP for model interpretation",
                "Data privacy regulations: GDPR, CCPA compliance",
                "Responsible AI development frameworks and guidelines",
            ]
        ),
    ]

    for title_text, body_lines in slides_data:
        # Use layout 1: Title and Content
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # Set title — regular weight (NOT bold)
        title_shape = slide.shapes.title
        title_shape.text = ""
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title_text
        run.font.size = Pt(32)
        run.font.bold = False  # Explicitly NOT bold
        run.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)  # Dark navy blue

        # Set body content
        body_shape = slide.placeholders[1]
        tf_body = body_shape.text_frame
        tf_body.clear()
        for i, line in enumerate(body_lines):
            if i == 0:
                p_body = tf_body.paragraphs[0]
            else:
                p_body = tf_body.add_paragraph()
            p_body.level = 0
            run_body = p_body.add_run()
            run_body.text = line
            run_body.font.size = Pt(18)
            run_body.font.bold = False
            run_body.font.color.rgb = RGBColor(0x26, 0x26, 0x26)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
