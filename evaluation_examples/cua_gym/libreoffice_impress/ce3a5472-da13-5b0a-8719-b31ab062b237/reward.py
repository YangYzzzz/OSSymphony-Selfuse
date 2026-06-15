"""
Reward Script: Apply bold formatting to all title textboxes in this presentation.
Task ID: osworld_impress_title_selective_formatting_002
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): Title placeholders on slides 1-3 are bold
  Component 2 (0.5): Title placeholders on slides 4-6 are bold
  Total: 1.0
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_selective_formatting_002'

TOTAL_SLIDES = 6
TITLE_TEXTS = [
    'Introduction to Machine Learning',
    'Supervised Learning Fundamentals',
    'Neural Networks and Deep Learning',
    'Natural Language Processing',
    'Model Evaluation and Validation',
    'Ethics and Fairness in AI',
]


def persist_app_state():
    """Attempt to save any open LibreOffice Impress document before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def get_title_bold_status(prs):
    """
    Return a list of (slide_index_1based, title_text, is_bold) for all slides.
    Checks the title placeholder (placeholder idx == 0) on each slide.
    A title is considered bold if ALL non-empty runs in ALL paragraphs are bold.
    Returns None for slides that have no title placeholder with runs.
    """
    results = []
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        title_shape = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape.is_placeholder:
                ph = shape.placeholder_format
                if ph.idx == 0:  # title placeholder index
                    title_shape = shape
                    break
        if title_shape is None:
            results.append((slide_num, None, None))
            continue

        # Collect all non-empty runs across all paragraphs
        all_runs = []
        for para in title_shape.text_frame.paragraphs:
            for run in para.runs:
                if (run.text or "").strip():
                    all_runs.append(run)

        if not all_runs:
            results.append((slide_num, title_shape.text, None))
            continue

        # All runs must be bold for the title to count as bold
        # Treat None (inherited) as not explicitly bold
        is_bold = all(run.font.bold is True for run in all_runs)
        results.append((slide_num, title_shape.text, is_bold))

    return results


def verify_task(file_path):
    """
    Verify that all 6 title placeholders are bold.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Gate: verify slide count is still 6
    if len(prs.slides) != TOTAL_SLIDES:
        print(f"FAIL: Expected {TOTAL_SLIDES} slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    title_statuses = get_title_bold_status(prs)

    # Component 1: Title placeholders on slides 1–3 are bold (0.5 points)
    try:
        slides_1_3 = [s for s in title_statuses if s[0] in (1, 2, 3)]
        all_bold_1_3 = all(s[2] is True for s in slides_1_3)
        if all_bold_1_3:
            print("PASS: Component 1 — Slides 1-3 title placeholders are bold (0.5 pts)")
            total_score += 0.5
        else:
            failed = [s for s in slides_1_3 if s[2] is not True]
            for s in failed:
                print(f"FAIL: Component 1 — Slide {s[0]} title '{s[1]}' is NOT bold (bold={s[2]})")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Title placeholders on slides 4–6 are bold (0.5 points)
    try:
        slides_4_6 = [s for s in title_statuses if s[0] in (4, 5, 6)]
        all_bold_4_6 = all(s[2] is True for s in slides_4_6)
        if all_bold_4_6:
            print("PASS: Component 2 — Slides 4-6 title placeholders are bold (0.5 pts)")
            total_score += 0.5
        else:
            failed = [s for s in slides_4_6 if s[2] is not True]
            for s in failed:
                print(f"FAIL: Component 2 — Slide {s[0]} title '{s[1]}' is NOT bold (bold={s[2]})")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point: test against canonical artifact path on the VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
