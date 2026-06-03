"""
Initial Setup: 6-slide product feature walkthrough presentation
Task ID: osworld_impress_title_selective_formatting_008
Domain: libreoffice_impress

Creates the initial-env artifact with:
- 6 slides, product feature walkthrough
- Slides 2, 3, 4, 5 all have titles in dark gray (no black, no underline)
- Slide 1: Title/intro slide
- Slide 6: Summary/closing slide
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_selective_formatting_008'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

# Dark gray color used for all slide titles (initial state)
DARK_GRAY = RGBColor(0x40, 0x40, 0x40)


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_title_format(title_shape, text, color, bold=True, underline=False, font_size=Pt(32)):
    """Set title text with specified formatting."""
    tf = title_shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.underline = underline


def add_content_text(slide, content_lines, left=Inches(0.8), top=Inches(2.2),
                     width=Inches(8.5), height=Inches(4.0)):
    """Add a text box with content lines."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title / Intro slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title1 = slide1.shapes.title
    title1.text = "NovaTech Product Suite"
    for para in title1.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1A, 0x53, 0x96)
    sub1 = slide1.placeholders[1]
    sub1.text = "2025 Feature Walkthrough Overview"
    for para in sub1.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # ---- Slide 2: Feature 1 - Dashboard Analytics ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    # Title text box
    title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(1.0))
    set_title_format(title_box2, "Dashboard Analytics", DARK_GRAY, bold=True, underline=False)
    # Content
    add_content_text(slide2, [
        "• Real-time KPI monitoring across all business units",
        "• Customizable widget layout with drag-and-drop interface",
        "• Automated daily, weekly, and monthly report generation",
        "• Integration with Salesforce, HubSpot, and Google Analytics",
        "• Role-based access controls for data visibility",
    ])
    # Decorative line
    line = slide2.shapes.add_shape(
        1, Inches(0.5), Inches(1.4), Inches(9.0), Emu(36000)
    )
    line.fill.background()
    line.line.color.rgb = RGBColor(0x1A, 0x53, 0x96)
    line.line.width = Emu(36000)

    # ---- Slide 3: Feature 2 - Collaboration Hub ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    title_box3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(1.0))
    set_title_format(title_box3, "Collaboration Hub", DARK_GRAY, bold=True, underline=False)
    add_content_text(slide3, [
        "• Shared workspaces for cross-functional teams",
        "• Version-controlled document editing with comment threads",
        "• Video conferencing with screen sharing and breakout rooms",
        "• Task assignment and milestone tracking with Gantt view",
        "• Slack and Microsoft Teams integration for instant notifications",
    ])
    line3 = slide3.shapes.add_shape(
        1, Inches(0.5), Inches(1.4), Inches(9.0), Emu(36000)
    )
    line3.fill.background()
    line3.line.color.rgb = RGBColor(0x1A, 0x53, 0x96)
    line3.line.width = Emu(36000)

    # ---- Slide 4: Feature 3 - Security & Compliance ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    title_box4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(1.0))
    set_title_format(title_box4, "Security & Compliance", DARK_GRAY, bold=True, underline=False)
    add_content_text(slide4, [
        "• End-to-end AES-256 encryption for all data at rest and in transit",
        "• SOC 2 Type II and ISO 27001 certified infrastructure",
        "• GDPR, CCPA, and HIPAA compliance built into data pipelines",
        "• Multi-factor authentication with biometric support",
        "• Automated vulnerability scanning and intrusion detection",
    ])
    line4 = slide4.shapes.add_shape(
        1, Inches(0.5), Inches(1.4), Inches(9.0), Emu(36000)
    )
    line4.fill.background()
    line4.line.color.rgb = RGBColor(0x1A, 0x53, 0x96)
    line4.line.width = Emu(36000)

    # ---- Slide 5: Feature 4 - AI-Powered Automation ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    title_box5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(1.0))
    set_title_format(title_box5, "AI-Powered Automation", DARK_GRAY, bold=True, underline=False)
    add_content_text(slide5, [
        "• Intelligent workflow automation with no-code builder",
        "• Predictive analytics for demand forecasting and churn prediction",
        "• NLP-based document classification and data extraction",
        "• Smart scheduling with calendar optimization algorithms",
        "• Anomaly detection in financial transactions and usage patterns",
    ])
    line5 = slide5.shapes.add_shape(
        1, Inches(0.5), Inches(1.4), Inches(9.0), Emu(36000)
    )
    line5.fill.background()
    line5.line.color.rgb = RGBColor(0x1A, 0x53, 0x96)
    line5.line.width = Emu(36000)

    # ---- Slide 6: Summary / Closing ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    title_box6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(1.0))
    set_title_format(title_box6, "Next Steps & Q&A", RGBColor(0x1A, 0x53, 0x96),
                     bold=True, underline=False, font_size=Pt(32))
    add_content_text(slide6, [
        "• Schedule a personalized demo with your account manager",
        "• Access free 30-day trial at novatech.io/trial",
        "• Join our community forum and knowledge base",
        "• Contact enterprise sales: enterprise@novatech.io",
        "",
        "Thank you for joining the NovaTech 2025 Feature Walkthrough!",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
