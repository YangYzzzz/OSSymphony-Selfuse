"""
Reward Script: Apply black color and underline to title text on slides 2, 3, and 5 only.
Task ID: osworld_impress_title_selective_formatting_008
Domain: libreoffice_impress
Scoring:
  - Slide 2 title color=black (000000): 0.15 pts
  - Slide 2 title underline=True:       0.10 pts
  - Slide 3 title color=black (000000): 0.15 pts
  - Slide 3 title underline=True:       0.10 pts
  - Slide 5 title color=black (000000): 0.15 pts
  - Slide 5 title underline=True:       0.10 pts
  - Slide 4 title remains dark gray (404040) and not underlined: 0.25 pts
  Total: 1.0

Context:
  Slides 2-6 use "TextBox 2" (bold text near top) as their visual title, not the placeholder.
  Initial state: all TextBox 2 titles have color=404040 (dark gray) and underline=False.
  Golden state:  slides 2, 3, 5 have color=000000 (black) and underline=True.
                 slide 4 is unchanged (404040, not underlined).
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_selective_formatting_008'


def get_title_textbox_run(slide):
    """
    Retrieve the first run of the title-like TextBox 2 on a slide.
    Returns the run or None if not found.
    """
    for shape in slide.shapes:
        if shape.name == 'TextBox 2' and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                runs = [r for r in para.runs if (r.text or '').strip()]
                if runs:
                    return runs[0]
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"CRITICAL: Expected at least 5 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]  # 0-indexed: slide 2
    slide3 = prs.slides[2]  # 0-indexed: slide 3
    slide4 = prs.slides[3]  # 0-indexed: slide 4
    slide5 = prs.slides[4]  # 0-indexed: slide 5

    # -------------------------------------------------------------------------
    # Component 1: Slide 2 title text color is black (000000) (0.15 pts)
    # -------------------------------------------------------------------------
    try:
        run = get_title_textbox_run(slide2)
        if run is not None and run.font.color.type is not None:
            color_str = str(run.font.color.rgb).upper()
            if color_str == '000000':
                print(f"PASS: Slide 2 title color is black (000000) (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Slide 2 title color expected 000000, found {color_str}")
        else:
            actual_color = str(run.font.color.rgb) if run and run.font.color.type is not None else 'inherited/None'
            print(f"FAIL: Slide 2 title color not set to black; run={run}, color={actual_color}")
    except Exception as e:
        print(f"ERROR: Component 1 (slide 2 color) — {e}")

    # -------------------------------------------------------------------------
    # Component 2: Slide 2 title text is underlined (0.10 pts)
    # -------------------------------------------------------------------------
    try:
        run = get_title_textbox_run(slide2)
        if run is not None and run.font.underline is True:
            print(f"PASS: Slide 2 title is underlined (0.10 pts)")
            total_score += 0.10
        else:
            actual_ul = run.font.underline if run is not None else 'run_not_found'
            print(f"FAIL: Slide 2 title underline expected True, found {actual_ul}")
    except Exception as e:
        print(f"ERROR: Component 2 (slide 2 underline) — {e}")

    # -------------------------------------------------------------------------
    # Component 3: Slide 3 title text color is black (000000) (0.15 pts)
    # -------------------------------------------------------------------------
    try:
        run = get_title_textbox_run(slide3)
        if run is not None and run.font.color.type is not None:
            color_str = str(run.font.color.rgb).upper()
            if color_str == '000000':
                print(f"PASS: Slide 3 title color is black (000000) (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Slide 3 title color expected 000000, found {color_str}")
        else:
            actual_color = str(run.font.color.rgb) if run and run.font.color.type is not None else 'inherited/None'
            print(f"FAIL: Slide 3 title color not set to black; color={actual_color}")
    except Exception as e:
        print(f"ERROR: Component 3 (slide 3 color) — {e}")

    # -------------------------------------------------------------------------
    # Component 4: Slide 3 title text is underlined (0.10 pts)
    # -------------------------------------------------------------------------
    try:
        run = get_title_textbox_run(slide3)
        if run is not None and run.font.underline is True:
            print(f"PASS: Slide 3 title is underlined (0.10 pts)")
            total_score += 0.10
        else:
            actual_ul = run.font.underline if run is not None else 'run_not_found'
            print(f"FAIL: Slide 3 title underline expected True, found {actual_ul}")
    except Exception as e:
        print(f"ERROR: Component 4 (slide 3 underline) — {e}")

    # -------------------------------------------------------------------------
    # Component 5: Slide 5 title text color is black (000000) (0.15 pts)
    # -------------------------------------------------------------------------
    try:
        run = get_title_textbox_run(slide5)
        if run is not None and run.font.color.type is not None:
            color_str = str(run.font.color.rgb).upper()
            if color_str == '000000':
                print(f"PASS: Slide 5 title color is black (000000) (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Slide 5 title color expected 000000, found {color_str}")
        else:
            actual_color = str(run.font.color.rgb) if run and run.font.color.type is not None else 'inherited/None'
            print(f"FAIL: Slide 5 title color not set to black; color={actual_color}")
    except Exception as e:
        print(f"ERROR: Component 5 (slide 5 color) — {e}")

    # -------------------------------------------------------------------------
    # Component 6: Slide 5 title text is underlined (0.10 pts)
    # -------------------------------------------------------------------------
    try:
        run = get_title_textbox_run(slide5)
        if run is not None and run.font.underline is True:
            print(f"PASS: Slide 5 title is underlined (0.10 pts)")
            total_score += 0.10
        else:
            actual_ul = run.font.underline if run is not None else 'run_not_found'
            print(f"FAIL: Slide 5 title underline expected True, found {actual_ul}")
    except Exception as e:
        print(f"ERROR: Component 6 (slide 5 underline) — {e}")

    # -------------------------------------------------------------------------
    # Component 7: Slide 4 title remains unchanged — dark gray (404040), NOT underlined (0.25 pts)
    # This verifies that the agent did NOT incorrectly modify slide 4.
    # In initial_env: color=404040, underline=False — this component FAILS initial
    # because it requires BOTH conditions to hold AND slides 2/3/5 to have passed.
    # Actually this component checks slide 4's untouched state which is already true
    # in initial. So we combine it: award points ONLY IF slides 2/3/5 were changed
    # (we use the score so far as a gate).
    # -------------------------------------------------------------------------
    try:
        run = get_title_textbox_run(slide4)
        if run is not None:
            try:
                color_str = str(run.font.color.rgb).upper() if run.font.color.type is not None else 'inherited'
            except Exception:
                color_str = 'err'
            underline = run.font.underline
            # Slide 4 should be unchanged: dark gray (404040) and NOT underlined
            color_ok = (color_str == '404040')
            underline_ok = (underline is False or underline is None)
            if color_ok and underline_ok:
                # Only award this point if the slides that SHOULD have changed (2, 3, 5)
                # were actually changed — i.e., at least some score was earned above.
                # This prevents initial_env from scoring 0.25 just because slide 4 is untouched.
                if total_score > 0.0:
                    print(f"PASS: Slide 4 title correctly unchanged — color={color_str}, underline={underline} (0.25 pts)")
                    total_score += 0.25
                else:
                    print(f"FAIL: Slide 4 title appears unchanged but no changes were made to slides 2/3/5 — "
                          f"component not awarded (slides 2/3/5 must be changed first)")
            else:
                print(f"FAIL: Slide 4 title was incorrectly modified — color={color_str} (expected 404040), "
                      f"underline={underline} (expected False/None)")
        else:
            print(f"FAIL: Could not find title TextBox 2 on slide 4")
    except Exception as e:
        print(f"ERROR: Component 7 (slide 4 unchanged) — {e}")

    final_score = round(min(total_score, 1.0), 4)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in a given env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
