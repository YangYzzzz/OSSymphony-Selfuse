"""
Reward Script: Process flowchart on slide 5 with diamond/rect/rounded-rect shapes and color coding
Task ID: impress_stu_049
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30) — Flowchart shapes with correct labels exist on slide 5
  Component 2 (0.25) — Correct shape types (rounded rect, rect, diamond)
  Component 3 (0.25) — Correct color coding (processes #D6EAF8, decisions #FEF9E7)
  Component 4 (0.20) — Connectors present and Yes/No decision labels
"""

import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_049'

# Expected flowchart labels
EXPECTED_LABELS = {'Start', 'Collect Data', 'Data Complete?', 'Analyze', 'Report', 'End'}

# Shape type constants from python-pptx MSO_SHAPE enum
# ROUNDED_RECTANGLE = 5, RECTANGLE = 1, DIAMOND = 4
EXPECTED_SHAPE_TYPES = {
    'Start': 5,           # ROUNDED_RECTANGLE
    'Collect Data': 1,    # RECTANGLE
    'Data Complete?': 4,  # DIAMOND
    'Analyze': 1,         # RECTANGLE
    'Report': 1,          # RECTANGLE
    'End': 5,             # ROUNDED_RECTANGLE
}

# Expected fill colors (hex uppercase, no '#')
PROCESS_COLOR = 'D6EAF8'   # light blue for rect processes
DECISION_COLOR = 'FEF9E7'  # light yellow for diamond decisions

EXPECTED_COLORS = {
    'Collect Data': PROCESS_COLOR,
    'Data Complete?': DECISION_COLOR,
    'Analyze': PROCESS_COLOR,
    'Report': PROCESS_COLOR,
}


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # 0-indexed, slide 5

    # Gather auto shapes on slide 5
    auto_shapes = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            label = shape.text.strip() if hasattr(shape, 'text') else ''
            auto_shapes.append({
                'name': shape.name,
                'label': label,
                'auto_shape_type': shape.auto_shape_type,
                'shape': shape,
            })

    # Build label -> shape info mapping
    label_map = {}
    for s in auto_shapes:
        if s['label']:
            label_map[s['label']] = s

    # -------------------------------------------------------
    # Component 1: Flowchart shapes with correct labels (0.30)
    # -------------------------------------------------------
    try:
        found_labels = set(label_map.keys()) & EXPECTED_LABELS
        match_ratio = len(found_labels) / len(EXPECTED_LABELS)
        if match_ratio >= 1.0:
            print(f"PASS: Component 1 — All 6 flowchart labels found: {sorted(found_labels)} (0.30 pts)")
            total_score += 0.30
        elif match_ratio > 0:
            missing = EXPECTED_LABELS - found_labels
            print(f"PARTIAL: Component 1 — {len(found_labels)}/6 labels found, missing: {missing}")
            if match_ratio > 0:
                total_score += round(0.30 * match_ratio, 2)
        else:
            print(f"FAIL: Component 1 — No expected flowchart labels found on slide 5. Found shapes: {[s['label'] for s in auto_shapes]}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # -------------------------------------------------------
    # Component 2: Correct shape types (0.25)
    # -------------------------------------------------------
    try:
        type_matches = 0
        type_total = 0
        for label, expected_type in EXPECTED_SHAPE_TYPES.items():
            if label in label_map:
                type_total += 1
                actual_type = label_map[label]['auto_shape_type']
                # auto_shape_type is an enum; compare int value
                actual_int = int(actual_type) if actual_type is not None else -1
                if actual_int == expected_type:
                    print(f"  PASS: '{label}' shape type = {actual_int} (expected {expected_type})")
                    type_matches += 1
                else:
                    print(f"  FAIL: '{label}' shape type = {actual_int}, expected {expected_type}")

        if type_total > 0:
            ratio = type_matches / len(EXPECTED_SHAPE_TYPES)
            pts = round(0.25 * ratio, 2)
            if ratio >= 1.0:
                print(f"PASS: Component 2 — All shape types correct (0.25 pts)")
                total_score += 0.25
            elif ratio > 0:
                print(f"PARTIAL: Component 2 — {type_matches}/{len(EXPECTED_SHAPE_TYPES)} shape types correct ({pts} pts)")
                total_score += pts
            else:
                print(f"FAIL: Component 2 — No shape types match")
        else:
            print(f"FAIL: Component 2 — No expected shapes found to check types")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # -------------------------------------------------------
    # Component 3: Correct color coding (0.25)
    # -------------------------------------------------------
    try:
        color_matches = 0
        color_total = len(EXPECTED_COLORS)
        for label, expected_hex in EXPECTED_COLORS.items():
            if label in label_map:
                shape_obj = label_map[label]['shape']
                try:
                    fill = shape_obj.fill
                    if fill.type is not None and fill.type == 1:  # SOLID fill
                        actual_rgb = str(fill.fore_color.rgb).upper()
                        if actual_rgb == expected_hex.upper():
                            print(f"  PASS: '{label}' fill color = {actual_rgb} (expected {expected_hex})")
                            color_matches += 1
                        else:
                            print(f"  FAIL: '{label}' fill color = {actual_rgb}, expected {expected_hex}")
                    else:
                        print(f"  FAIL: '{label}' fill type is {fill.type}, not SOLID (1)")
                except Exception as e:
                    print(f"  FAIL: '{label}' could not read fill color: {e}")
            else:
                print(f"  SKIP: '{label}' not found on slide")

        if color_matches > 0:
            ratio = color_matches / color_total
            pts = round(0.25 * ratio, 2)
            if ratio >= 1.0:
                print(f"PASS: Component 3 — All process/decision colors correct (0.25 pts)")
                total_score += 0.25
            else:
                print(f"PARTIAL: Component 3 — {color_matches}/{color_total} colors correct")
                if color_matches > 0:
                    total_score += round(0.25 * (color_matches / color_total), 2)
        else:
            print(f"FAIL: Component 3 — No colors match expected values")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # -------------------------------------------------------
    # Component 4: Connectors and Yes/No labels (0.20)
    # -------------------------------------------------------
    try:
        # Count connectors (LINE type shapes)
        connectors = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.LINE]
        num_connectors = len(connectors)

        # Check for Yes/No text labels (TextBox type)
        text_boxes = []
        for s in slide.shapes:
            if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and hasattr(s, 'text'):
                text_boxes.append(s.text.strip())

        has_yes = any('Yes' in t for t in text_boxes)
        has_no = any('No' in t for t in text_boxes)

        # Need at least 5 connectors (Start->Collect, Collect->Decision, Decision->Analyze, Analyze->Report, Report->End, plus loop-back)
        # and both Yes/No labels
        sub_score = 0.0

        if num_connectors >= 5:
            print(f"  PASS: {num_connectors} connectors found (need >= 5)")
            sub_score += 0.10
        elif num_connectors >= 3:
            print(f"  PARTIAL: {num_connectors} connectors found (need >= 5)")
            sub_score += 0.05
        else:
            print(f"  FAIL: Only {num_connectors} connectors found (need >= 5)")

        if has_yes and has_no:
            print(f"  PASS: Both 'Yes' and 'No' decision labels found")
            sub_score += 0.10
        elif has_yes or has_no:
            found_label = 'Yes' if has_yes else 'No'
            print(f"  PARTIAL: Only '{found_label}' label found, missing the other")
            sub_score += 0.05
        else:
            print(f"  FAIL: Neither 'Yes' nor 'No' decision labels found")

        if sub_score > 0:
            total_score += sub_score
            print(f"{'PASS' if sub_score >= 0.20 else 'PARTIAL'}: Component 4 — Connectors & labels ({sub_score} pts)")
        else:
            print(f"FAIL: Component 4 — No connectors or labels found")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
