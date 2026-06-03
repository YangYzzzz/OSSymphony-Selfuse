"""
Initial Setup: Create a 7-slide Research Workflow presentation with empty slide 5
Task ID: impress_stu_049
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_049'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_body(slide, title_text, body_lines):
    """Helper to set title and body content on a slide."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    # Find content placeholder (index 1 typically)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.paragraphs[0].text = body_lines[0] if body_lines else ""
            for line in body_lines[1:]:
                p = tf.add_paragraph()
                p.text = line
                p.level = 0
            break


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Research Workflow"
    slide1.placeholders[1].text = "Data Collection & Analysis Framework\nQ2 2025 Update"

    # --- Slide 2: Project Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide2, "Project Overview", [
        "This presentation outlines the research workflow adopted by the Analytics team.",
        "Key objectives include streamlining data collection pipelines,",
        "improving quality assurance checks, and automating report generation.",
        "The framework has been tested across 3 pilot projects since January 2025.",
    ])

    # --- Slide 3: Team Members ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide3, "Team Members", [
        "Dr. Sarah Chen — Principal Investigator",
        "Marcus Johnson — Data Engineer",
        "Priya Patel — Statistical Analyst",
        "James O'Brien — Research Coordinator",
        "Aiko Tanaka — Quality Assurance Lead",
    ])

    # --- Slide 4: Data Sources ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide4, "Data Sources", [
        "Internal CRM database (PostgreSQL) — 2.4M records",
        "Survey responses via Qualtrics — 12,500 respondents",
        "Public API feeds (Census Bureau, BLS) — updated monthly",
        "Partner data sharing agreements — 3 external organizations",
        "Web scraping pipeline — 45 target domains, weekly cadence",
    ])

    # --- Slide 5: Research Process (EMPTY — task target) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add only a title text box at the top
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Research Process"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    # Body is intentionally empty — the agent will add the flowchart here

    # --- Slide 6: Timeline ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide6, "Project Timeline", [
        "Phase 1 (Jan–Mar): Infrastructure setup and pilot testing",
        "Phase 2 (Apr–Jun): Full-scale data collection rollout",
        "Phase 3 (Jul–Sep): Analysis and model development",
        "Phase 4 (Oct–Dec): Report generation and stakeholder review",
    ])

    # --- Slide 7: Next Steps ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide7, "Next Steps", [
        "Finalize data sharing agreements with remaining partners",
        "Deploy automated quality checks on ingestion pipeline",
        "Schedule mid-project review with advisory board (June 15)",
        "Begin drafting interim report for Q3 distribution",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
