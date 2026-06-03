"""
Initial Setup: Create a 4-slide presentation with slide 3 titled 'Essay Rubric' but no table content.
Task ID: impress_teach_036
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_036'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Assignment Rubric Collection"
    slide1.placeholders[1].text = "English Literature — Fall 2025\nProfessor Sarah Mitchell"

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Grading Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "This presentation contains rubrics for all major assignments."
    p2 = body2.add_paragraph()
    p2.text = "Each rubric uses a 4-point scale from Excellent to Needs Work."
    p3 = body2.add_paragraph()
    p3.text = "Students should review rubrics before submitting assignments."
    p4 = body2.add_paragraph()
    p4.text = "Rubrics are aligned with course learning objectives."

    # --- Slide 3: Essay Rubric (title only, NO table) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Essay Rubric"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 4: Participation Rubric ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Participation Rubric"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Class participation is evaluated weekly based on:"
    items = [
        "Quality of contributions to class discussions",
        "Preparedness and completion of assigned readings",
        "Respectful engagement with peers' ideas",
        "Attendance and punctuality"
    ]
    for item in items:
        p = body4.add_paragraph()
        p.text = item
        p.level = 1

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
