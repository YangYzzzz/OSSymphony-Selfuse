"""
Reward Script: Insert a grading rubric table on slide 3 with headers and alternating row colors
Task ID: impress_teach_036
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): Table exists on slide 3 with 6 rows x 5 columns
  Component 2 (0.3): Header row has correct text values
  Component 3 (0.4): Data rows (1-5) have alternating fill colors #E8F5E9 and #FFFFFF
"""

import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_036'

EXPECTED_HEADERS = ['Criteria', 'Excellent (4)', 'Good (3)', 'Satisfactory (2)', 'Needs Work (1)']

# Alternating colors for data rows (rows 1-5): odd data rows = E8F5E9, even data rows = FFFFFF
# Row 1 (data row 1, odd) = E8F5E9
# Row 2 (data row 2, even) = FFFFFF
# Row 3 (data row 3, odd) = E8F5E9
# Row 4 (data row 4, even) = FFFFFF
# Row 5 (data row 5, odd) = E8F5E9
EXPECTED_ROW_COLORS = {
    1: 'E8F5E9',
    2: 'FFFFFF',
    3: 'E8F5E9',
    4: 'FFFFFF',
    5: 'E8F5E9',
}


def get_cell_fill_color(cell):
    """Extract the solid fill color hex from a table cell, or None."""
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is not None:
        solidFill = tcPr.find(qn('a:solidFill'))
        if solidFill is not None:
            srgbClr = solidFill.find(qn('a:srgbClr'))
            if srgbClr is not None:
                return srgbClr.get('val')
    return None


def find_table_on_slide(slide):
    """Find the first table shape on a slide, or None."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return shape.table
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]  # 0-indexed

    # Component 1: Table exists on slide 3 with 6 rows x 5 columns (0.3 points)
    try:
        table = find_table_on_slide(slide3)
        if table is None:
            print("FAIL: Component 1 -- No table found on slide 3")
        else:
            rows = len(table.rows)
            cols = len(table.columns)
            if rows == 6 and cols == 5:
                print(f"PASS: Component 1 -- Table found on slide 3 with {rows}x{cols} dimensions (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 1 -- Table dimensions are {rows}x{cols}, expected 6x5")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Header row has correct text values (0.3 points)
    try:
        table = find_table_on_slide(slide3)
        if table is not None and len(table.rows) >= 1 and len(table.columns) >= 5:
            actual_headers = [table.cell(0, c).text.strip() for c in range(5)]
            matching = sum(1 for a, e in zip(actual_headers, EXPECTED_HEADERS) if a == e)
            if matching == 5:
                print(f"PASS: Component 2 -- All 5 headers match exactly (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 2 -- {matching}/5 headers match. Actual: {actual_headers}")
        else:
            print("FAIL: Component 2 -- Table not found or too small to check headers")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Data rows have alternating fill colors (0.4 points)
    # Each correct row color earns 0.08 pts (5 rows x 0.08 = 0.4)
    try:
        table = find_table_on_slide(slide3)
        if table is not None and len(table.rows) >= 6:
            color_score = 0.0
            for row_idx in range(1, 6):
                expected_color = EXPECTED_ROW_COLORS[row_idx]
                # Check color of first cell in the row as representative
                # But verify all cells in the row have same color for robustness
                row_colors = []
                for col_idx in range(min(5, len(table.columns))):
                    c = get_cell_fill_color(table.cell(row_idx, col_idx))
                    row_colors.append(c)

                # Check if all cells in row have expected color
                all_match = all(
                    (c is not None and c.upper() == expected_color.upper())
                    for c in row_colors
                )
                if all_match:
                    print(f"PASS: Row {row_idx} fill color = #{expected_color} (0.08 pts)")
                    color_score += 0.08
                else:
                    # Partial: check if at least first cell matches
                    first_color = row_colors[0]
                    if first_color is not None and first_color.upper() == expected_color.upper():
                        print(f"PARTIAL: Row {row_idx} first cell matches #{expected_color} but not all cells")
                    else:
                        print(f"FAIL: Row {row_idx} expected #{expected_color}, got colors: {row_colors}")

            if color_score > 0:
                total_score += color_score
            if color_score >= 0.39:
                print(f"PASS: Component 3 -- All data row colors correct ({color_score:.2f} pts)")
            else:
                print(f"PARTIAL: Component 3 -- Some row colors correct ({color_score:.2f}/0.40 pts)")
        else:
            print("FAIL: Component 3 -- Table not found or too few rows")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
