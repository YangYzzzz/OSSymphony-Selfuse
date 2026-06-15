"""
Initial Setup: Insert a pie chart on slide 3
Task ID: impress_tm_051
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_051'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=None, alignment=None):
    """Helper to add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tf


def add_bullet_paragraph(text_frame, text, level=0, font_size=16, color=None):
    """Add a bulleted paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    run = p.runs[0]
    run.font.size = Pt(font_size)
    if color:
        run.font.color.rgb = color
    return p


def create_initial():
    prs = Presentation()

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Annual Budget Overview 2025"
    slide1.placeholders[1].text = "Financial Planning & Resource Allocation"

    # ---- Slide 2: Revenue Summary ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Revenue Summary"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Total revenue for FY2025: $12.8 million"
    add_bullet_paragraph(body2, "Q1: $2.9M (+8% YoY)", level=1)
    add_bullet_paragraph(body2, "Q2: $3.2M (+12% YoY)", level=1)
    add_bullet_paragraph(body2, "Q3: $3.4M (+15% YoY)", level=1)
    add_bullet_paragraph(body2, "Q4: $3.3M (projected)", level=1)

    # ---- Slide 3: Budget Allocation (TITLE ONLY - no chart) ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title text box manually
    add_text_box(
        slide3,
        Inches(0.5), Inches(0.3), Inches(9.0), Inches(1.0),
        "Budget Allocation",
        font_size=32,
        bold=True,
        color=RGBColor(0x1F, 0x49, 0x7D),
        alignment=PP_ALIGN.LEFT,
    )

    # ---- Slide 4: Timeline ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Implementation Timeline"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Key milestones for budget execution:"
    add_bullet_paragraph(body4, "January - March: Initial allocation and team setup", level=1)
    add_bullet_paragraph(body4, "April - June: Mid-year review and adjustments", level=1)
    add_bullet_paragraph(body4, "July - September: Performance evaluation", level=1)
    add_bullet_paragraph(body4, "October - December: Final reconciliation", level=1)

    # ---- Slide 5: Next Steps ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Next Steps"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Action items for the leadership team:"
    add_bullet_paragraph(body5, "Review department-level budget proposals by Feb 15", level=1)
    add_bullet_paragraph(body5, "Schedule quarterly budget review meetings", level=1)
    add_bullet_paragraph(body5, "Finalize vendor contracts for engineering tools", level=1)
    add_bullet_paragraph(body5, "Submit compliance documentation to Finance", level=1)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
