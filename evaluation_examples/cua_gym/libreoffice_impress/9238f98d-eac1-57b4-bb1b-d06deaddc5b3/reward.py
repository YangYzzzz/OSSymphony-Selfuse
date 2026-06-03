"""
Reward Script: Insert a pie chart on slide 3 with budget allocation data
Task ID: impress_tm_051
Domain: libreoffice_impress
Scoring:
  - Component 1: Chart exists on slide 3 and is a pie chart (0.35 pts)
  - Component 2: Chart categories match Marketing, Engineering, Sales (0.30 pts)
  - Component 3: Chart values match 35%, 40%, 25% (0.20 pts)
  - Component 4: Chart has a legend (0.15 pts)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_051'


def persist_app_state(domain: str):
    """Save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.enum.chart import XL_CHART_TYPE
    except ImportError as e:
        print(f"CRITICAL: Missing library: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]  # 0-indexed

    # Find chart shape(s) on slide 3
    chart_shapes = []
    for shape in slide3.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            chart_shapes.append(shape)

    # Component 1: Chart exists on slide 3 and is a pie chart (0.35 points)
    pie_chart_types = {
        XL_CHART_TYPE.PIE,
        XL_CHART_TYPE.PIE_EXPLODED,
        XL_CHART_TYPE.THREE_D_PIE,
        XL_CHART_TYPE.THREE_D_PIE_EXPLODED,
    }
    try:
        if len(chart_shapes) == 0:
            print("FAIL: Component 1 — No chart found on slide 3")
        else:
            chart = chart_shapes[0].chart
            if chart.chart_type in pie_chart_types:
                print(f"PASS: Component 1 — Pie chart found on slide 3 (type={chart.chart_type}) (0.35 pts)")
                total_score += 0.35
            else:
                print(f"FAIL: Component 1 — Chart found but not a pie chart (type={chart.chart_type})")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Chart categories are Marketing, Engineering, Sales (0.30 points)
    try:
        if len(chart_shapes) > 0:
            chart = chart_shapes[0].chart
            categories = list(chart.plots[0].categories)
            expected_cats = ['Marketing', 'Engineering', 'Sales']
            # Case-insensitive comparison
            cats_lower = [c.strip().lower() for c in categories]
            expected_lower = [c.lower() for c in expected_cats]
            if cats_lower == expected_lower:
                print(f"PASS: Component 2 — Categories match: {categories} (0.30 pts)")
                total_score += 0.30
            elif set(cats_lower) == set(expected_lower):
                # Categories present but in different order — partial credit
                print(f"PARTIAL: Component 2 — Categories present but in different order: {categories} (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 2 — Expected categories {expected_cats}, found {categories}")
        else:
            print("FAIL: Component 2 — No chart to check categories")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Chart values match 35, 40, 25 (0.20 points)
    try:
        if len(chart_shapes) > 0:
            chart = chart_shapes[0].chart
            series = chart.plots[0].series[0]
            values = list(series.values)
            categories = list(chart.plots[0].categories)
            # Build a category→value map for order-independent comparison
            cat_val = {}
            for cat, val in zip(categories, values):
                cat_val[cat.strip().lower()] = val

            expected_map = {'marketing': 35.0, 'engineering': 40.0, 'sales': 25.0}
            mismatches = []
            for cat, exp_val in expected_map.items():
                actual = cat_val.get(cat)
                if actual is None:
                    mismatches.append(f"Missing category '{cat}' in data")
                elif abs(actual - exp_val) > 0.5:
                    mismatches.append(f"{cat} expected {exp_val}, found {actual}")

            if len(mismatches) == 0 and len(cat_val) >= 3:
                print(f"PASS: Component 3 — Values match: {dict(zip(categories, values))} (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 3 — {'; '.join(mismatches) if mismatches else 'insufficient categories'}")
        else:
            print("FAIL: Component 3 — No chart to check values")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Chart has a legend (0.15 points)
    try:
        if len(chart_shapes) > 0:
            chart = chart_shapes[0].chart
            if chart.has_legend:
                print(f"PASS: Component 4 — Chart has legend (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 4 — Chart does not have a legend")
        else:
            print("FAIL: Component 4 — No chart to check legend")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
