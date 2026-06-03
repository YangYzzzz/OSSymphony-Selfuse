"""
Initial Setup: Create a 6-slide presentation with text on slide 5, no animations.
Task ID: impress_ma_073
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_073'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Word Reveal Techniques"
    slide1.placeholders[1].text = "Creative Presentation Methods for 2025"

    # --- Slide 2: Overview (Title + Content) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Introduction to word reveal effects"
    p2a = body2.add_paragraph()
    p2a.text = "Benefits of progressive text display"
    p2a.level = 0
    p2b = body2.add_paragraph()
    p2b.text = "Case studies from marketing campaigns"
    p2b.level = 0
    p2c = body2.add_paragraph()
    p2c.text = "Implementation best practices"
    p2c.level = 0
    p2d = body2.add_paragraph()
    p2d.text = "Q&A and next steps"
    p2d.level = 0

    # --- Slide 3: Background on Audience Engagement ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Why Progressive Reveal Matters"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Studies show that revealing text gradually increases audience retention by 34%"
    p3a = body3.add_paragraph()
    p3a.text = "Viewers engage more deeply when information arrives in digestible segments"
    p3a.level = 0
    p3b = body3.add_paragraph()
    p3b.text = "Presenters can control pacing and emphasize key points"
    p3b.level = 0

    # --- Slide 4: Case Study ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Case Study: Nextera Product Launch"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Nextera used word-by-word reveals for their Q3 product launch"
    p4a = body4.add_paragraph()
    p4a.text = "Audience engagement scores increased from 72% to 91%"
    p4a.level = 1
    p4b = body4.add_paragraph()
    p4b.text = "Average attention span during presentation grew by 8 minutes"
    p4b.level = 1
    p4c = body4.add_paragraph()
    p4c.text = "Post-event survey showed 96% recall of key messaging"
    p4c.level = 1

    # --- Slide 5: The Key Statement (target for animation task) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide5.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p5 = tf.paragraphs[0]
    p5.text = "Innovation drives our success forward every day"
    p5.alignment = PP_ALIGN.CENTER
    run = p5.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x7A)

    # --- Slide 6: Summary ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Key Takeaways"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Progressive text reveal keeps audiences focused"
    p6a = body6.add_paragraph()
    p6a.text = "Word-by-word animation is ideal for impactful statements"
    p6a.level = 0
    p6b = body6.add_paragraph()
    p6b.text = "Combine with transitions for maximum effect"
    p6b.level = 0
    p6c = body6.add_paragraph()
    p6c.text = "Practice timing to match your speaking cadence"
    p6c.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
