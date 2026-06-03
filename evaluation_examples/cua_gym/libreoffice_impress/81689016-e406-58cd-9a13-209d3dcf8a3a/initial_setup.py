"""
Initial Setup: Photo Gallery presentation with master slide title at top-left, left-aligned.
Task ID: impress_ma_019
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_019'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # Default slide dimensions: 10in x 7.5in
    # Ensure standard dimensions
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Modify the master slide title placeholder ---
    # Access the first slide master
    master = prs.slide_masters[0]

    # Find the title placeholder on the master and position it at top-left
    for ph in master.placeholders:
        if ph.placeholder_format.idx == 0:  # Title placeholder
            # Position at top-left area
            ph.left = Inches(0.5)
            ph.top = Inches(0.3)
            ph.width = Inches(5.0)
            ph.height = Inches(1.2)
            # Set text alignment to LEFT
            for para in ph.text_frame.paragraphs:
                para.alignment = PP_ALIGN.LEFT
            break

    # Slide content data for a photo gallery presentation
    slide_data = [
        {
            'layout_idx': 0,  # Title Slide
            'title': 'Urban Landscapes Photography',
            'subtitle': 'A Visual Journey Through Metropolitan Architecture\nBy Elena Vasquez | Spring 2025 Collection'
        },
        {
            'layout_idx': 1,  # Title + Content
            'title': 'Exhibition Overview',
            'body': 'This collection features 48 photographs captured across\n12 major cities during a 6-month creative residency.\n\nKey themes include:\n- Geometric patterns in modern architecture\n- Light and shadow interplay at golden hour\n- Human presence in urban environments\n- Reflections and symmetry in glass facades'
        },
        {
            'layout_idx': 1,
            'title': 'New York City Series',
            'body': 'Brooklyn Bridge at Dawn - Shot at 5:42 AM\nMidtown Reflections - Glass towers on 5th Avenue\nCentral Park Framing - Natural meets architectural\nSubway Motion - Long exposure at Grand Central\n\nCamera: Nikon Z9 | Lens: 24-70mm f/2.8'
        },
        {
            'layout_idx': 1,
            'title': 'Tokyo Nightscapes',
            'body': 'Shibuya Crossing - 2,500 pedestrians per cycle\nShinjuku Neon District - Blade Runner aesthetic\nTokyo Tower from Roppongi Hills\nSenso-ji Temple at Twilight\n\nISO range: 3200-12800 | Aperture: f/1.4-f/4'
        },
        {
            'layout_idx': 1,
            'title': 'Barcelona Architecture',
            'body': 'Sagrada Familia Interior Light Study\nCasa Batllo Facade Detail - Gaudi mosaic work\nGothic Quarter Alleyways - Medieval geometry\nW Hotel Sail Structure - Contemporary curves\n\nAll prints available in 24x36 and 30x40 formats'
        },
        {
            'layout_idx': 1,
            'title': 'Dubai Skyline Collection',
            'body': 'Burj Khalifa Vanishing Point - Looking straight up\nPalm Jumeirah Aerial Perspective\nDubai Marina at Blue Hour\nDesert Meets City - Sand dunes with skyline backdrop\n\nDrone photography: DJI Mavic 3 Pro'
        },
        {
            'layout_idx': 1,
            'title': 'London Through the Lens',
            'body': 'The Shard in Morning Fog\nTower Bridge Mechanics - Engineering as art\nSt Pauls Dome - Symmetry study\nBrick Lane Street Art - Urban canvas culture\n\nFilm stock: Kodak Portra 400 (select images)'
        },
        {
            'layout_idx': 1,
            'title': 'Technical Approach',
            'body': 'Equipment Used:\n- Primary: Nikon Z9 with 24-70mm f/2.8 S\n- Wide: Nikon 14-24mm f/2.8\n- Tele: Nikon 70-200mm f/2.8 VR S\n- Drone: DJI Mavic 3 Pro\n- Tripod: Gitzo Systematic Series 3\n\nPost-processing: Adobe Lightroom + Photoshop\nColor grading: Custom LUT profiles per city'
        },
        {
            'layout_idx': 1,
            'title': 'Print Specifications',
            'body': 'Paper: Hahnemuhle Photo Rag Baryta 315gsm\nPrinter: Epson SureColor P9570\nColor profile: Adobe RGB (1998)\nEdition sizes: 25 signed + 5 AP per image\n\nPricing:\n  24x36 framed: $1,850\n  30x40 framed: $2,400\n  40x60 museum mount: $4,200'
        },
        {
            'layout_idx': 1,
            'title': 'Exhibition Schedule',
            'body': 'Upcoming Venues:\n- Aperture Gallery, NYC: June 12 - July 28, 2025\n- Foam Museum, Amsterdam: August 15 - October 5, 2025\n- C/O Berlin: November 1 - December 20, 2025\n- SFMOMA Pop-up: January 2026\n\nPrivate viewings available by appointment'
        },
        {
            'layout_idx': 1,
            'title': 'Artist Statement',
            'body': 'Cities are living organisms. Their architecture forms a\nskeleton, their inhabitants the lifeblood. I seek the\nmoments where structure and humanity intersect - \nwhere rigid geometry yields to organic movement.\n\nEach image in this collection represents a conversation\nbetween the built environment and the people who\ninhabit it. The camera is simply the translator.'
        },
        {
            'layout_idx': 1,
            'title': 'Contact & Inquiries',
            'body': 'Elena Vasquez Photography\nStudio: 142 West 26th Street, Floor 8, New York, NY 10001\n\nEmail: elena@vasquezphoto.com\nWeb: www.vasquezphoto.com\nInstagram: @elenavasquezphoto\n\nRepresentation: Meridian Fine Art Gallery\nAgent: Thomas Hartwell | thomas@meridiangallery.com'
        },
    ]

    for sdata in slide_data:
        layout_idx = sdata['layout_idx']
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = sdata['title']
            # Ensure title text is left-aligned (matching the master)
            for para in slide.shapes.title.text_frame.paragraphs:
                para.alignment = PP_ALIGN.LEFT

        # Set body/subtitle
        if 'subtitle' in sdata and 1 in slide.placeholders:
            slide.placeholders[1].text = sdata['subtitle']
        elif 'body' in sdata and 1 in slide.placeholders:
            slide.placeholders[1].text = sdata['body']

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
