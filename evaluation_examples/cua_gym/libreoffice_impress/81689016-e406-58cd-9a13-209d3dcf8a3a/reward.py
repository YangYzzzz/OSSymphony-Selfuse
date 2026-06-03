"""
Reward Script: Move master slide title placeholder to vertically centered with center text alignment
Task ID: impress_ma_019
Domain: libreoffice_impress
Scoring:
  Component 1 (0.35): Master slide title placeholder vertically centered (~3.0in top)
  Component 2 (0.25): Master slide title text alignment is CENTER
  Component 3 (0.25): Majority of slide-level titles vertically repositioned to ~3.0in
  Component 4 (0.15): Majority of slide-level titles have CENTER alignment
"""

import os

from pptx import Presentation
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_019'

# Golden target: title placeholder top ~3.0 inches on 7.5in slide
# 3.0 inches = 2743200 EMU
TARGET_TOP_EMU = 2743200
SLIDE_HEIGHT_EMU = 6858000  # 7.5 inches

# Tolerance: 10% of slide height (~0.75in) for vertical centering
VERTICAL_TOLERANCE_EMU = int(SLIDE_HEIGHT_EMU * 0.10)


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Master slide title placeholder is vertically centered (0.35 points)
    # Initial: top=274320 (0.30in). Golden: top=2743200 (3.00in).
    # The title should be near the vertical center of the slide.
    try:
        master = prs.slide_masters[0]
        title_ph = None
        for ph in master.placeholders:
            if ph.placeholder_format.idx == 0:  # title placeholder
                title_ph = ph
                break

        if title_ph is None:
            print("FAIL: Component 1 -- No title placeholder found on master slide")
        else:
            master_top = title_ph.top
            ph_height = title_ph.height
            # Vertical center of placeholder
            ph_center = master_top + ph_height // 2
            # Expected center of slide
            slide_center = SLIDE_HEIGHT_EMU // 2

            diff = abs(ph_center - slide_center)
            print(f"  Master title top: {master_top} EMU ({master_top/914400:.2f}in)")
            print(f"  PH center: {ph_center} EMU, Slide center: {slide_center} EMU, Diff: {diff} EMU")

            if diff <= VERTICAL_TOLERANCE_EMU:
                print(f"PASS: Component 1 -- Master title vertically centered (0.35 pts)")
                total_score += 0.35
            else:
                print(f"FAIL: Component 1 -- Master title not vertically centered. Top={master_top/914400:.2f}in, expected ~3.00in")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Master slide title text alignment is CENTER (0.25 points)
    # Initial: LEFT. Golden: CENTER.
    try:
        master = prs.slide_masters[0]
        title_ph = None
        for ph in master.placeholders:
            if ph.placeholder_format.idx == 0:
                title_ph = ph
                break

        if title_ph is None:
            print("FAIL: Component 2 -- No title placeholder found on master slide")
        elif not title_ph.has_text_frame:
            print("FAIL: Component 2 -- Title placeholder has no text frame")
        else:
            # Check alignment of first paragraph
            alignment = title_ph.text_frame.paragraphs[0].alignment
            if alignment == PP_ALIGN.CENTER:
                print(f"PASS: Component 2 -- Master title alignment is CENTER (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 2 -- Master title alignment is {alignment}, expected CENTER")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Majority of slide-level titles repositioned to ~3.0in top (0.25 points)
    # Initial: most at 0.30in. Golden: all at 3.00in.
    try:
        centered_count = 0
        total_with_title = 0
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if (hasattr(shape, 'placeholder_format') and
                        shape.placeholder_format is not None and
                        shape.placeholder_format.idx == 0):
                    total_with_title += 1
                    shape_center = shape.top + shape.height // 2
                    slide_center_y = SLIDE_HEIGHT_EMU // 2
                    if abs(shape_center - slide_center_y) <= VERTICAL_TOLERANCE_EMU:
                        centered_count += 1
                    break

        if total_with_title == 0:
            print("FAIL: Component 3 -- No title placeholders found on any slide")
        else:
            ratio = centered_count / total_with_title
            print(f"  Slides with centered title: {centered_count}/{total_with_title} ({ratio:.0%})")
            if ratio >= 0.75:
                print(f"PASS: Component 3 -- Majority of slide titles vertically centered (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 3 -- Only {centered_count}/{total_with_title} slides have centered titles")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Majority of slide-level titles have CENTER alignment (0.15 points)
    # Initial: all LEFT. Golden: all CENTER.
    try:
        center_aligned_count = 0
        total_with_title = 0
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if (hasattr(shape, 'placeholder_format') and
                        shape.placeholder_format is not None and
                        shape.placeholder_format.idx == 0):
                    total_with_title += 1
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.alignment == PP_ALIGN.CENTER:
                                center_aligned_count += 1
                            break  # check first paragraph only
                    break

        if total_with_title == 0:
            print("FAIL: Component 4 -- No title placeholders found on any slide")
        else:
            ratio = center_aligned_count / total_with_title
            print(f"  Slides with center-aligned title: {center_aligned_count}/{total_with_title} ({ratio:.0%})")
            if ratio >= 0.75:
                print(f"PASS: Component 4 -- Majority of slide titles center-aligned (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 4 -- Only {center_aligned_count}/{total_with_title} slides have center alignment")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
