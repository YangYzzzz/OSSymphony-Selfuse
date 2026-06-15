"""
Reward Script: Board Meeting Presentation from Board_Data.xlsx
Task ID: impress_wf_089
Domain: libreoffice_impress
Scoring:
  C1: File exists on Desktop (0.05)
  C2: Exactly 12 slides (0.10)
  C3: Slide 1 title text (0.10)
  C4: Fade transitions on all 12 slides (0.15)
  C5: Navy #0D47A1 and Gold #C9A959 colors used (0.15)
  C6: Slide 5 has P&L table (0.10)
  C7: Slide 8 competitive landscape quadrant labels (0.10)
  C8: Slide 9 risks table (0.10)
  C9: Slide 10 M&A pipeline table (0.05)
  C10: Slide 11 has 3 resolution cards with vote boxes (0.10)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_089'


def get_all_text(slide):
    """Extract all text from a slide, including nested shapes."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return texts


def check_fade_transitions(pptx_path, num_slides):
    """Check that all slides have fade transitions via XML parsing."""
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    fade_count = 0
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            for i in range(1, num_slides + 1):
                fname = f'ppt/slides/slide{i}.xml'
                try:
                    with zf.open(fname) as f:
                        root = ET.parse(f).getroot()
                        tr = root.find(f'.//{{{ns_p}}}transition')
                        if tr is not None:
                            fade = tr.find(f'.//{{{ns_p}}}fade')
                            if fade is not None:
                                fade_count += 1
                except KeyError:
                    pass
    except Exception as e:
        print(f"ERROR: Cannot parse ZIP for transitions: {e}")
    return fade_count


def check_colors_used(prs):
    """Check if navy (#0D47A1) and gold (#C9A959) are used."""
    navy_found = False
    gold_found = False

    for slide in prs.slides:
        # Check background
        try:
            bg = slide.background.fill
            if bg.type == 1:
                rgb = str(bg.fore_color.rgb).upper()
                if rgb == '0D47A1':
                    navy_found = True
                if rgb == 'C9A959':
                    gold_found = True
        except Exception:
            pass

        # Check shape fills and text colors
        for shape in slide.shapes:
            # Shape fill
            if shape.shape_type == 1:  # AUTO_SHAPE
                try:
                    fill = shape.fill
                    if fill.type == 1:
                        rgb = str(fill.fore_color.rgb).upper()
                        if rgb == '0D47A1':
                            navy_found = True
                        if rgb == 'C9A959':
                            gold_found = True
                except Exception:
                    pass

            # Text color
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color.type is not None:
                                rgb = str(run.font.color.rgb).upper()
                                if rgb == '0D47A1':
                                    navy_found = True
                                if rgb == 'C9A959':
                                    gold_found = True
                        except Exception:
                            pass

            if navy_found and gold_found:
                return True, True

    return navy_found, gold_found


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    total_score = 0.0

    # Component 1: File exists on Desktop (0.05 points)
    # This is task-introduced: initial_env has no Board_Meeting.pptx
    if not os.path.exists(file_path):
        print(f"FAIL: Component 1 — File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load presentation {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    print(f"PASS: Component 1 — File exists and is valid pptx (0.05 pts)")
    total_score += 0.05

    num_slides = len(prs.slides)

    # Component 2: Exactly 12 slides (0.10 points)
    try:
        if num_slides == 12:
            print(f"PASS: Component 2 — Slide count is 12 (0.10 pts)")
            total_score += 0.10
        elif num_slides >= 10:
            partial = 0.05
            print(f"PARTIAL: Component 2 — Slide count is {num_slides}, expected 12 ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — Slide count is {num_slides}, expected 12")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 1 title text contains key phrases (0.10 points)
    try:
        if num_slides >= 1:
            slide1_texts = get_all_text(prs.slides[0])
            all_text_lower = " ".join(slide1_texts).lower()
            has_board = "board" in all_text_lower and ("director" in all_text_lower or "meeting" in all_text_lower)
            has_q3 = "q3" in all_text_lower and "2024" in all_text_lower

            if has_board and has_q3:
                print(f"PASS: Component 3 — Slide 1 has 'Board of Directors Meeting' and 'Q3 2024' (0.10 pts)")
                total_score += 0.10
            elif has_board or has_q3:
                print(f"PARTIAL: Component 3 — Slide 1 has partial title match (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 3 — Slide 1 text: {slide1_texts[:3]}")
        else:
            print("FAIL: Component 3 — No slides found")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Fade transitions on all 12 slides (0.15 points)
    try:
        fade_count = check_fade_transitions(file_path, num_slides)
        if fade_count == num_slides and num_slides >= 12:
            print(f"PASS: Component 4 — All {num_slides} slides have Fade transitions (0.15 pts)")
            total_score += 0.15
        elif fade_count >= 6:
            partial = round(0.15 * (fade_count / max(num_slides, 12)), 2)
            print(f"PARTIAL: Component 4 — {fade_count}/{num_slides} slides have Fade transitions ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — Only {fade_count}/{num_slides} slides have Fade transitions")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Navy (#0D47A1) and Gold (#C9A959) colors used (0.15 points)
    try:
        navy, gold = check_colors_used(prs)
        if navy and gold:
            print(f"PASS: Component 5 — Both navy and gold colors found (0.15 pts)")
            total_score += 0.15
        elif navy or gold:
            print(f"PARTIAL: Component 5 — Navy={navy}, Gold={gold} (0.075 pts)")
            total_score += 0.075
        else:
            print(f"FAIL: Component 5 — Neither navy nor gold colors found")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Slide 5 has P&L table (0.10 points)
    try:
        if num_slides >= 5:
            slide5 = prs.slides[4]
            slide5_texts = get_all_text(slide5)
            all_text_s5 = " ".join(slide5_texts).lower()

            has_table = False
            table_rows = 0
            for shape in slide5.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    has_table = True
                    table_rows = len(shape.table.rows)
                    break

            has_pnl_title = ("p&l" in all_text_s5 or "profit" in all_text_s5 or
                             "loss" in all_text_s5 or "p & l" in all_text_s5)

            if has_table and has_pnl_title:
                print(f"PASS: Component 6 — Slide 5 has P&L table ({table_rows} rows) (0.10 pts)")
                total_score += 0.10
            elif has_table:
                print(f"PARTIAL: Component 6 — Slide 5 has table but no P&L title (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 6 — Slide 5 has no table. Texts: {slide5_texts[:3]}")
        else:
            print("FAIL: Component 6 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Slide 8 competitive landscape with quadrant labels (0.10 points)
    try:
        if num_slides >= 8:
            slide8 = prs.slides[7]
            slide8_texts = get_all_text(slide8)
            all_text_s8 = " ".join(slide8_texts).lower()

            has_competitive = ("competitive" in all_text_s8 or "landscape" in all_text_s8)
            # Check for quadrant labels
            quadrant_labels = ["leader", "challenger", "niche", "laggard"]
            found_labels = sum(1 for label in quadrant_labels if label in all_text_s8)

            if has_competitive and found_labels >= 3:
                print(f"PASS: Component 7 — Slide 8 has competitive landscape with {found_labels}/4 quadrant labels (0.10 pts)")
                total_score += 0.10
            elif has_competitive and found_labels >= 1:
                print(f"PARTIAL: Component 7 — Slide 8 has competitive title with {found_labels}/4 labels (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 7 — Slide 8 missing competitive landscape content. Texts: {slide8_texts[:5]}")
        else:
            print("FAIL: Component 7 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    # Component 8: Slide 9 has risks table (0.10 points)
    try:
        if num_slides >= 9:
            slide9 = prs.slides[8]
            slide9_texts = get_all_text(slide9)
            all_text_s9 = " ".join(slide9_texts).lower()

            has_table = False
            table_cols = 0
            for shape in slide9.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    has_table = True
                    table_cols = len(shape.table.columns)
                    # Check headers
                    headers = [shape.table.cell(0, c).text.lower() for c in range(table_cols)]
                    break

            has_risk_title = "risk" in all_text_s9

            if has_table and has_risk_title:
                print(f"PASS: Component 8 — Slide 9 has risks table ({table_cols} columns) (0.10 pts)")
                total_score += 0.10
            elif has_table:
                print(f"PARTIAL: Component 8 — Slide 9 has table but no risk title (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 8 — Slide 9 has no table")
        else:
            print("FAIL: Component 8 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 8 — {e}")

    # Component 9: Slide 10 M&A pipeline table (0.05 points)
    try:
        if num_slides >= 10:
            slide10 = prs.slides[9]
            slide10_texts = get_all_text(slide10)
            all_text_s10 = " ".join(slide10_texts).lower()

            has_table = False
            has_correct_headers = False
            for shape in slide10.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    has_table = True
                    t = shape.table
                    headers = [t.cell(0, c).text.lower() for c in range(len(t.columns))]
                    # Check for Target, Stage, Valuation
                    if ("target" in headers and "stage" in headers and
                            "valuation" in headers):
                        has_correct_headers = True
                    break

            has_ma_title = "m&a" in all_text_s10 or "pipeline" in all_text_s10

            if has_table and has_correct_headers and has_ma_title:
                print(f"PASS: Component 9 — Slide 10 has M&A pipeline table with correct headers (0.05 pts)")
                total_score += 0.05
            elif has_table and has_ma_title:
                print(f"PARTIAL: Component 9 — Slide 10 has table and M&A title but wrong headers (0.025 pts)")
                total_score += 0.025
            else:
                print(f"FAIL: Component 9 — Slide 10 missing M&A pipeline table")
        else:
            print("FAIL: Component 9 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 9 — {e}")

    # Component 10: Slide 11 has 3 resolution cards with vote boxes (0.10 points)
    try:
        if num_slides >= 11:
            slide11 = prs.slides[10]
            slide11_texts = get_all_text(slide11)
            all_text_s11 = " ".join(slide11_texts).lower()

            # Count resolution references
            resolution_count = all_text_s11.count("resolution")
            # Count vote option keywords
            approve_count = all_text_s11.count("approve")
            reject_count = all_text_s11.count("reject")
            abstain_count = all_text_s11.count("abstain")

            has_governance = ("governance" in all_text_s11 or "resolution" in all_text_s11)
            has_3_resolutions = resolution_count >= 3
            has_vote_boxes = (approve_count >= 3 and reject_count >= 3 and abstain_count >= 3)

            if has_governance and has_3_resolutions and has_vote_boxes:
                print(f"PASS: Component 10 — Slide 11 has {resolution_count} resolutions with vote boxes (0.10 pts)")
                total_score += 0.10
            elif has_governance and has_3_resolutions:
                print(f"PARTIAL: Component 10 — Slide 11 has resolutions but incomplete vote boxes (0.05 pts)")
                total_score += 0.05
            elif has_governance:
                print(f"PARTIAL: Component 10 — Slide 11 has governance title only (0.03 pts)")
                total_score += 0.03
            else:
                print(f"FAIL: Component 10 — Slide 11 missing governance content")
        else:
            print("FAIL: Component 10 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 10 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/Desktop/Board_Meeting.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
