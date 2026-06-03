"""
Initial Setup: Add Fly In animation to chart on slide 5
Task ID: impress_sales_032
Domain: libreoffice_impress

Creates an 8-slide sales presentation. Slide 5 has a title with a Fade
entrance animation and a bar chart with NO animation.
"""

import os
import shlex
import subprocess
import time
import zipfile
import shutil
from io import BytesIO
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_032'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=None, alignment=None):
    """Helper to add a text box with formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullet_slide(slide, title_text, bullets):
    """Add title and bullet points to a slide."""
    # Title
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                 title_text, font_size=28, bold=True,
                 color=RGBColor(0x1A, 0x3C, 0x6E))
    # Bullets
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title Slide ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    bg = slide1.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
    add_text_box(slide1, Inches(1), Inches(2), Inches(11), Inches(1.5),
                 "Q3 2025 Sales Performance Review", font_size=36, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, Inches(1), Inches(3.8), Inches(11), Inches(1),
                 "Global Sales Division  |  Presented by Sarah Chen, VP Sales",
                 font_size=18, color=RGBColor(0xCC, 0xCC, 0xCC),
                 alignment=PP_ALIGN.CENTER)

    # ── Slide 2: Executive Summary ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    add_bullet_slide(slide2, "Executive Summary", [
        "Total revenue reached $12.8M, a 15% increase over Q2 2025",
        "APAC region contributed 38% of total revenue, up from 31%",
        "Enterprise segment grew 22% driven by new cloud partnerships",
        "Customer retention rate improved to 94.2% from 91.7%",
        "New product line 'Apex Suite' exceeded launch targets by 40%",
    ])

    # ── Slide 3: Regional Breakdown (table) ──
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide3, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                 "Regional Revenue Breakdown", font_size=28, bold=True,
                 color=RGBColor(0x1A, 0x3C, 0x6E))
    tbl_shape = slide3.shapes.add_table(
        6, 4, Inches(0.8), Inches(1.5), Inches(11), Inches(4))
    tbl = tbl_shape.table
    headers = ["Region", "Q3 Revenue ($M)", "Q2 Revenue ($M)", "Growth (%)"]
    data_rows = [
        ["North America", "$4.23", "$3.92", "+7.9%"],
        ["APAC", "$4.86", "$3.85", "+26.2%"],
        ["EMEA", "$2.51", "$2.38", "+5.5%"],
        ["Latin America", "$0.82", "$0.73", "+12.3%"],
        ["Middle East & Africa", "$0.38", "$0.30", "+26.7%"],
    ]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
    for r, row_data in enumerate(data_rows, 1):
        for c, val in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(13)

    # ── Slide 4: Product Line Analysis ──
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_bullet_slide(slide4, "Product Line Analysis", [
        "Apex Suite: $3.2M revenue (25% of total) — exceeded forecast by 40%",
        "Enterprise Cloud: $4.1M revenue — steady 8% growth, high margin",
        "SMB Platform: $2.9M revenue — new pricing model boosted conversions",
        "Professional Services: $1.8M — consulting demand up in APAC",
        "Legacy Products: $0.8M — planned decline, migration path active",
    ])

    # ── Slide 5: Revenue by Quarter (BAR CHART + TITLE) ──
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    # Title shape
    title_shape = add_text_box(
        slide5, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
        "Revenue by Quarter", font_size=28, bold=True,
        color=RGBColor(0x1A, 0x3C, 0x6E))

    # Bar chart
    chart_data = CategoryChartData()
    chart_data.categories = ['Q1 2025', 'Q2 2025', 'Q3 2025']
    chart_data.add_series('North America', (3650000, 3920000, 4230000))
    chart_data.add_series('APAC', (3100000, 3850000, 4860000))
    chart_data.add_series('EMEA', (2200000, 2380000, 2510000))
    chart_data.add_series('LATAM', (650000, 730000, 820000))

    chart_shape = slide5.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.8),
        Inches(11), Inches(5), chart_data)

    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.include_in_layout = False

    # ── Slide 6: Market Trends ──
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_bullet_slide(slide6, "Market Trends & Competitive Landscape", [
        "AI-driven analytics tools saw 45% market growth across all segments",
        "Competitors increased cloud investment — Zenith Corp launched rival suite",
        "Customer demand shifting toward integrated platform solutions",
        "Regulatory changes in EMEA creating new compliance opportunities",
        "Partner ecosystem expanded by 18 new certified integrators",
    ])

    # ── Slide 7: Customer Acquisition ──
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_bullet_slide(slide7, "Customer Acquisition & Retention", [
        "247 new enterprise accounts acquired in Q3 (vs 198 in Q2)",
        "Average deal size increased to $52K from $47K",
        "Sales cycle shortened by 12 days through new demo automation",
        "Churn rate decreased to 5.8% from 8.3% year-over-year",
        "Net Promoter Score reached 72, highest in company history",
    ])

    # ── Slide 8: Next Steps ──
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    add_bullet_slide(slide8, "Next Steps & Action Items", [
        "Expand APAC sales team by 15 headcount in Q4 2025",
        "Launch Apex Suite v2.0 with AI analytics module by November",
        "Pilot new partner certification program in North America",
        "Migrate remaining legacy customers to Enterprise Cloud by Q1 2026",
        "Target $14.5M revenue for Q4 with focus on enterprise upsells",
    ])

    # Save the base file first
    prs.save(OUTPUT)
    print(f'Base file saved: {OUTPUT}')

    # Now inject Fade animation on the title shape of slide 5 via XML
    # We need to identify the shape IDs first
    # Reload to get shape IDs
    prs2 = Presentation(OUTPUT)
    slide5_reload = prs2.slides[4]  # 0-indexed
    title_sp_id = None
    chart_sp_id = None
    for shape in slide5_reload.shapes:
        if shape.has_text_frame and "Revenue by Quarter" in shape.text:
            title_sp_id = shape.shape_id
        if shape.has_chart:
            chart_sp_id = shape.shape_id
    print(f'Slide 5 title shape ID: {title_sp_id}, chart shape ID: {chart_sp_id}')

    if title_sp_id is None:
        print("ERROR: Could not find title shape on slide 5")
        return

    # Inject animation XML into slide5.xml via ZIP manipulation
    inject_fade_animation(OUTPUT, title_sp_id)
    print(f'Fade animation injected on title (shape {title_sp_id})')
    print(f'Chart (shape {chart_sp_id}) has NO animation')

    # Launch GUI
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


def inject_fade_animation(pptx_path, title_shape_id):
    """Inject a Fade entrance animation on the title shape of slide 5."""
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

    # Register namespaces to avoid ns0/ns1 prefixes
    ET.register_namespace('a', ns_a)
    ET.register_namespace('p', ns_p)
    ET.register_namespace('r', ns_r)

    tmp_path = pptx_path + '.tmp'
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        with zipfile.ZipFile(tmp_path, 'w') as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == 'ppt/slides/slide5.xml':
                    data = _add_fade_to_slide_xml(data, title_shape_id)
                zout.writestr(item, data)

    shutil.move(tmp_path, pptx_path)


def _add_fade_to_slide_xml(xml_bytes, title_shape_id):
    """Add timing/animation XML for a Fade entrance on the title shape."""
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    root = ET.fromstring(xml_bytes)

    # Build the animation timing XML
    # The structure: p:timing > p:tnLst > p:par (main seq) > p:cTn > p:childTnLst > p:par (click) > ...
    timing_xml = f'''<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:tnLst>
        <p:par>
          <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
            <p:childTnLst>
              <p:seq concurrent="1" nextAc="seek">
                <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                  <p:childTnLst>
                    <p:par>
                      <p:cTn id="3" fill="hold">
                        <p:stCondLst>
                          <p:cond delay="0"/>
                        </p:stCondLst>
                        <p:childTnLst>
                          <p:par>
                            <p:cTn id="4" fill="hold">
                              <p:stCondLst>
                                <p:cond delay="0"/>
                              </p:stCondLst>
                              <p:childTnLst>
                                <p:animEffect transition="in" filter="fade">
                                  <p:cBhvr>
                                    <p:cTn id="5" dur="500" fill="hold"/>
                                    <p:tgtEl>
                                      <p:spTgt spid="{title_shape_id}"/>
                                    </p:tgtEl>
                                  </p:cBhvr>
                                </p:animEffect>
                              </p:childTnLst>
                            </p:cTn>
                          </p:par>
                        </p:childTnLst>
                      </p:cTn>
                    </p:par>
                  </p:childTnLst>
                </p:cTn>
                <p:prevCondLst>
                  <p:cond evt="onPrev" delay="0">
                    <p:tgtEl>
                      <p:sldTgt/>
                    </p:tgtEl>
                  </p:cond>
                </p:prevCondLst>
                <p:nextCondLst>
                  <p:cond evt="onClick" delay="0">
                    <p:tgtEl>
                      <p:sldTgt/>
                    </p:tgtEl>
                  </p:cond>
                </p:nextCondLst>
              </p:seq>
            </p:childTnLst>
          </p:cTn>
        </p:par>
      </p:tnLst>
    </p:timing>'''

    timing_el = ET.fromstring(timing_xml)

    # Remove existing timing if any
    existing = root.find(f'{{{ns_p}}}timing')
    if existing is not None:
        root.remove(existing)

    root.append(timing_el)

    return ET.tostring(root, xml_declaration=True, encoding='UTF-8')


create_initial()
