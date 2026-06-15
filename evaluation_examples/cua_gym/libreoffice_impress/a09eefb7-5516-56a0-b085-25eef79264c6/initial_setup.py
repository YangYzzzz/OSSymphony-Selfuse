"""
Initial Setup: Embed audio clips as background music for slide ranges
Task ID: impress_gf5_028
Domain: libreoffice_impress

Creates a 12-slide art portfolio presentation and 3 MP3 audio stubs.
No audio is embedded - that is the agent's task.
"""

import os
import shlex
import struct
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_028'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
AUDIO_DIR = f'{WORKDIR}/audio'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_minimal_mp3(filepath, duration_sec=5):
    """Create a minimal valid MP3 file (silent) for testing purposes."""
    os.makedirs(os.path.dirname(filepath), exist_ok=True)
    # Create a minimal MP3 with MPEG audio frame headers (silent frames)
    # MPEG1 Layer3, 128kbps, 44100Hz, stereo
    frame_header = b'\xff\xfb\x90\x00'
    # Each frame is 417 bytes at 128kbps/44100Hz
    frame_size = 417
    padding = b'\x00' * (frame_size - len(frame_header))
    frame = frame_header + padding
    # ~38 frames per second at 128kbps
    num_frames = 38 * duration_sec
    with open(filepath, 'wb') as f:
        # ID3v2 header (minimal)
        f.write(b'ID3')
        f.write(b'\x03\x00')  # version 2.3
        f.write(b'\x00')      # flags
        f.write(b'\x00\x00\x00\x00')  # size=0 (no tags)
        for _ in range(num_frames):
            f.write(frame)
    print(f'Created audio file: {filepath} ({duration_sec}s)')


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color palette for the art portfolio
    dark_bg = RGBColor(0x1A, 0x1A, 0x2E)
    accent_gold = RGBColor(0xD4, 0xAF, 0x37)
    light_text = RGBColor(0xF5, 0xF5, 0xF0)
    subtle_gray = RGBColor(0x8A, 0x8A, 0x8A)

    def set_bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_text(slide, text, left, top, width, height,
                 font_name="Georgia", font_size=24, color=light_text,
                 bold=False, alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = alignment
        run = p.runs[0]
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        return txBox

    # --- Slide 1: Title ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide1, dark_bg)
    add_text(slide1, "Echoes of Light", Inches(1.5), Inches(2), Inches(10), Inches(1.5),
             font_size=54, color=accent_gold, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide1, "A Contemporary Art Portfolio by Elena Vasquez",
             Inches(1.5), Inches(4), Inches(10), Inches(1),
             font_size=28, alignment=PP_ALIGN.CENTER)
    add_text(slide1, "Gallery Opening  |  March 2026  |  The Meridian Arts Center",
             Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
             font_size=16, color=subtle_gray, alignment=PP_ALIGN.CENTER)

    # --- Slide 2: Artist Statement ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide2, dark_bg)
    add_text(slide2, "Artist Statement", Inches(1), Inches(0.5), Inches(8), Inches(1),
             font_size=36, color=accent_gold, bold=True)
    statement = (
        "My work explores the interplay between natural light and urban geometry. "
        "Through layered mixed media, I seek to capture those fleeting moments when "
        "sunlight transforms ordinary spaces into something transcendent. Each piece "
        "invites the viewer to pause and rediscover the beauty in their everyday surroundings."
    )
    add_text(slide2, statement, Inches(1), Inches(2), Inches(10), Inches(4),
             font_size=20, color=light_text)

    # --- Slide 3: Exhibition Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide3, dark_bg)
    add_text(slide3, "Exhibition Overview", Inches(1), Inches(0.5), Inches(8), Inches(1),
             font_size=36, color=accent_gold, bold=True)
    overview_items = [
        "Series I: Urban Dawn — Large-scale oil on canvas (4 works)",
        "Series II: Reflected Geometries — Mixed media on panel (3 works)",
        "Series III: Twilight Passages — Digital prints on aluminum (3 works)",
        "Installation: Luminous Corridor — Site-specific light installation",
    ]
    txBox = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(overview_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.space_after = Pt(12)
        run = p.runs[0]
        run.font.name = "Georgia"
        run.font.size = Pt(18)
        run.font.color.rgb = light_text

    # --- Slides 4-6: Series I - Urban Dawn ---
    series1_works = [
        ("Morning on 5th Avenue", "Oil on canvas, 60\" x 84\"\n\nCapturing the first rays of sunlight as they cut between skyscrapers, casting long golden shadows across the awakening city streets."),
        ("Rooftop Solstice", "Oil on canvas, 48\" x 72\"\n\nA view from above as summer solstice light bathes a rooftop garden in warm amber, blurring the line between nature and architecture."),
        ("Bridge at Dawn", "Oil on canvas, 54\" x 78\"\n\nThe Brooklyn Bridge emerges from morning mist, its cables catching light like strings of a celestial harp."),
    ]
    for title, desc in series1_works:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, dark_bg)
        add_text(slide, title, Inches(1), Inches(0.5), Inches(10), Inches(1),
                 font_size=32, color=accent_gold, bold=True)
        add_text(slide, desc, Inches(1), Inches(2), Inches(5), Inches(4),
                 font_size=16, color=light_text)

    # --- Slides 7-8: Series II - Reflected Geometries ---
    series2_works = [
        ("Glass Canyon", "Mixed media on panel, 36\" x 48\"\n\nLayers of torn maps, metallic leaf, and oil paint create an abstract landscape of reflections between glass towers."),
        ("Puddle Universe", "Mixed media on panel, 30\" x 40\"\n\nA rain puddle becomes a portal, reflecting a fractured skyline in watercolor, ink, and collaged photography."),
    ]
    for title, desc in series2_works:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, dark_bg)
        add_text(slide, title, Inches(1), Inches(0.5), Inches(10), Inches(1),
                 font_size=32, color=accent_gold, bold=True)
        add_text(slide, desc, Inches(1), Inches(2), Inches(5), Inches(4),
                 font_size=16, color=light_text)

    # --- Slide 9: Series III - Twilight Passages ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide9, dark_bg)
    add_text(slide9, "Twilight Passages", Inches(1), Inches(0.5), Inches(10), Inches(1),
             font_size=36, color=accent_gold, bold=True)
    add_text(slide9, (
        "Digital prints on brushed aluminum, 40\" x 60\" each\n\n"
        "This series uses long-exposure photography composited with procedural "
        "geometry to create luminous pathways through familiar urban corridors. "
        "The aluminum substrate adds a reflective quality that shifts with the "
        "viewer's movement, echoing the ephemeral nature of twilight itself."
    ), Inches(1), Inches(2), Inches(10), Inches(4),
             font_size=18, color=light_text)

    # --- Slide 10: Installation - Luminous Corridor ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide10, dark_bg)
    add_text(slide10, "Luminous Corridor", Inches(1), Inches(0.5), Inches(10), Inches(1),
             font_size=36, color=accent_gold, bold=True)
    add_text(slide10, "Site-specific light installation", Inches(1), Inches(1.5), Inches(10), Inches(0.5),
             font_size=20, color=subtle_gray)
    add_text(slide10, (
        "A 40-foot corridor of programmable LED panels responds to viewer "
        "movement, recreating the shifting light patterns observed during golden "
        "hour across different cities. Sensors track visitors and adjust color "
        "temperature and intensity in real time, creating a unique experience "
        "with each passage."
    ), Inches(1), Inches(2.5), Inches(10), Inches(4),
             font_size=18, color=light_text)

    # --- Slide 11: Acknowledgments ---
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide11, dark_bg)
    add_text(slide11, "Acknowledgments", Inches(1), Inches(0.5), Inches(10), Inches(1),
             font_size=36, color=accent_gold, bold=True)
    ack_text = (
        "With deep gratitude to:\n\n"
        "The Meridian Arts Center and Director Priya Sundaram\n"
        "Whitfield Gallery — representation since 2019\n"
        "The National Endowment for the Arts — Grant #2025-VA-0847\n"
        "Studio assistants: Tomoko Murakami, David Osei, Lena Brandt\n"
        "My partner, Rafael, for unwavering support"
    )
    add_text(slide11, ack_text, Inches(1), Inches(2), Inches(10), Inches(4),
             font_size=18, color=light_text)

    # --- Slide 12: Contact / Closing ---
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide12, dark_bg)
    add_text(slide12, "Thank You", Inches(1.5), Inches(2), Inches(10), Inches(1.5),
             font_size=54, color=accent_gold, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide12, (
        "Elena Vasquez\n"
        "elena@vasquezstudio.art\n"
        "www.vasquezstudio.art\n"
        "@elenavasquezart"
    ), Inches(1.5), Inches(4), Inches(10), Inches(2),
             font_size=20, color=light_text, alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Slide count: {len(prs.slides)}')

    # Create audio files
    create_minimal_mp3(f'{AUDIO_DIR}/intro.mp3', duration_sec=30)
    create_minimal_mp3(f'{AUDIO_DIR}/main.mp3', duration_sec=60)
    create_minimal_mp3(f'{AUDIO_DIR}/outro.mp3', duration_sec=20)

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
