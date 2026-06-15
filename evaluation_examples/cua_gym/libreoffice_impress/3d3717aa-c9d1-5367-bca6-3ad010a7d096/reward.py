"""
Reward Script: Network Topology Diagram on Slide 5
Task ID: impress_ps_031
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20) - Cloud shape labeled 'Internet' exists on slide 5
  Component 2 (0.15) - Rectangle labeled 'Firewall' exists on slide 5
  Component 3 (0.15) - Rectangle labeled 'Core Switch' exists on slide 5
  Component 4 (0.20) - Three rectangles labeled 'Web Server 1/2/3' exist on slide 5
  Component 5 (0.15) - Line connectors present (at least 4)
  Component 6 (0.15) - Correct vertical layout (Internet top, servers bottom)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_031'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # 0-indexed, slide 5
    shapes = list(slide.shapes)

    # Collect auto shapes and connectors on slide 5 (excluding pre-existing title/textbox)
    auto_shapes = []
    connectors = []
    for s in shapes:
        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            auto_shapes.append(s)
        elif int(s.shape_type) == 9:  # LINE connector
            connectors.append(s)

    # Helper: find shape with matching text (case-insensitive, stripped)
    def find_shape_by_text(shape_list, target_text):
        for s in shape_list:
            if hasattr(s, 'text') and s.text.strip().lower() == target_text.strip().lower():
                return s
        return None

    # Component 1: Cloud shape labeled 'Internet' (0.20 points)
    try:
        internet_shape = find_shape_by_text(auto_shapes, 'Internet')
        if internet_shape is not None:
            # Verify it's a cloud shape
            try:
                is_cloud = (internet_shape.auto_shape_type is not None and
                            int(internet_shape.auto_shape_type) == 179)  # CLOUD
            except Exception:
                is_cloud = False

            if is_cloud:
                print(f"PASS: Component 1 - Cloud shape labeled 'Internet' found (0.20 pts)")
                total_score += 0.20
            elif internet_shape is not None:
                # It has the right text but not a cloud shape - partial
                print(f"PARTIAL: Component 1 - Shape labeled 'Internet' found but not cloud type (0.10 pts)")
                total_score += 0.10
        else:
            print(f"FAIL: Component 1 - No shape labeled 'Internet' found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Rectangle labeled 'Firewall' (0.15 points)
    try:
        firewall_shape = find_shape_by_text(auto_shapes, 'Firewall')
        if firewall_shape is not None:
            try:
                is_rect = (firewall_shape.auto_shape_type is not None and
                           int(firewall_shape.auto_shape_type) == 1)  # RECTANGLE
            except Exception:
                is_rect = False

            if is_rect:
                print(f"PASS: Component 2 - Rectangle labeled 'Firewall' found (0.15 pts)")
                total_score += 0.15
            elif firewall_shape is not None:
                print(f"PARTIAL: Component 2 - Shape labeled 'Firewall' found but not rectangle (0.08 pts)")
                total_score += 0.08
        else:
            print(f"FAIL: Component 2 - No shape labeled 'Firewall' found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Rectangle labeled 'Core Switch' (0.15 points)
    try:
        coreswitch_shape = find_shape_by_text(auto_shapes, 'Core Switch')
        if coreswitch_shape is not None:
            try:
                is_rect = (coreswitch_shape.auto_shape_type is not None and
                           int(coreswitch_shape.auto_shape_type) == 1)
            except Exception:
                is_rect = False

            if is_rect:
                print(f"PASS: Component 3 - Rectangle labeled 'Core Switch' found (0.15 pts)")
                total_score += 0.15
            elif coreswitch_shape is not None:
                print(f"PARTIAL: Component 3 - Shape labeled 'Core Switch' found but not rectangle (0.08 pts)")
                total_score += 0.08
        else:
            print(f"FAIL: Component 3 - No shape labeled 'Core Switch' found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Three rectangles labeled 'Web Server 1', 'Web Server 2', 'Web Server 3' (0.20 pts)
    try:
        ws_labels = ['Web Server 1', 'Web Server 2', 'Web Server 3']
        ws_found = 0
        for label in ws_labels:
            ws_shape = find_shape_by_text(auto_shapes, label)
            if ws_shape is not None:
                ws_found += 1
                print(f"  Found: '{label}'")
            else:
                print(f"  Missing: '{label}'")

        if ws_found == 3:
            print(f"PASS: Component 4 - All 3 web server shapes found (0.20 pts)")
            total_score += 0.20
        elif ws_found > 0:
            partial = round(0.20 * ws_found / 3, 2)
            print(f"PARTIAL: Component 4 - {ws_found}/3 web server shapes found ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 - No web server shapes found")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Line connectors present - at least 4 (0.15 points)
    try:
        num_connectors = len(connectors)
        if num_connectors >= 4:
            print(f"PASS: Component 5 - {num_connectors} line connectors found (>= 4) (0.15 pts)")
            total_score += 0.15
        elif num_connectors > 0:
            partial = round(0.15 * num_connectors / 4, 2)
            print(f"PARTIAL: Component 5 - {num_connectors}/4 connectors found ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 - No line connectors found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    # Component 6: Correct vertical layout - Internet at top, servers at bottom (0.15 points)
    try:
        layout_ok = (internet_shape is not None and firewall_shape is not None and coreswitch_shape is not None)
        issues = []

        # Internet should be above Firewall
        if internet_shape and firewall_shape:
            if internet_shape.top >= firewall_shape.top:
                layout_ok = False
                issues.append("Internet not above Firewall")
        elif not layout_ok:
            issues.append("Missing Internet or Firewall shape for layout check")

        # Firewall should be above Core Switch
        if firewall_shape and coreswitch_shape:
            if firewall_shape.top >= coreswitch_shape.top:
                layout_ok = False
                issues.append("Firewall not above Core Switch")
        elif not layout_ok:
            issues.append("Missing Firewall or Core Switch for layout check")

        # Core Switch should be above web servers
        if coreswitch_shape:
            for label in ws_labels:
                ws_shape = find_shape_by_text(auto_shapes, label)
                if ws_shape and ws_shape.top <= coreswitch_shape.top:
                    layout_ok = False
                    issues.append(f"'{label}' not below Core Switch")

        if layout_ok and internet_shape and firewall_shape and coreswitch_shape:
            print(f"PASS: Component 6 - Correct vertical layout (Internet > Firewall > Core Switch > Servers) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 6 - Layout issues: {'; '.join(issues)}")
    except Exception as e:
        print(f"ERROR: Component 6 - {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress (save any unsaved changes)
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
