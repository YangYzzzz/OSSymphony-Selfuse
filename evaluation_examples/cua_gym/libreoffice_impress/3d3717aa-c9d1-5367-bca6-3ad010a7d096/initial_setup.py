"""
Initial Setup: Create Tech_Architecture.pptx with 8 slides.
Slide 5 has title 'Network Topology' and empty content area.
Task ID: impress_ps_031
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_031'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Tech Architecture Overview"
    slide1.placeholders[1].text = "Infrastructure & Systems Design\nQ2 2025 Review"

    # --- Slide 2: System Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "System Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Cloud-native microservices architecture"
    p = body2.add_paragraph()
    p.text = "Deployed across 3 availability zones"
    p.level = 1
    p = body2.add_paragraph()
    p.text = "99.97% uptime SLA target"
    p.level = 1
    p = body2.add_paragraph()
    p.text = "Auto-scaling with Kubernetes orchestration"
    p.level = 1

    # --- Slide 3: Technology Stack ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Technology Stack"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Frontend: React 18, TypeScript, Next.js"
    for item in [
        "Backend: Python 3.11, FastAPI, gRPC",
        "Database: PostgreSQL 15, Redis 7, Elasticsearch 8",
        "Messaging: Apache Kafka, RabbitMQ",
        "CI/CD: GitHub Actions, ArgoCD, Terraform",
        "Monitoring: Prometheus, Grafana, PagerDuty",
    ]:
        p = body3.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 4: Data Flow Architecture ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Data Flow Architecture"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Event-driven data pipeline processing 2.4M events/day"
    for item in [
        "Ingestion: Kafka topics partitioned by region",
        "Processing: Spark Structured Streaming with exactly-once semantics",
        "Storage: Delta Lake on S3 with time-travel enabled",
        "Serving: Low-latency Redis cache + PostgreSQL read replicas",
    ]:
        p = body4.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 5: Network Topology (EMPTY - this is the task target) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add only a title text box at the top
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Network Topology"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x34, 0x40)
    # NO shapes, NO connectors, NO diagram elements — empty content area

    # --- Slide 6: Security Architecture ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Security Architecture"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Zero-trust security model with defense in depth"
    for item in [
        "mTLS for all service-to-service communication",
        "OAuth 2.0 + OIDC for user authentication",
        "HashiCorp Vault for secrets management",
        "WAF and DDoS protection at edge layer",
        "SOC 2 Type II compliant infrastructure",
    ]:
        p = body6.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 7: Disaster Recovery ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Disaster Recovery Plan"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "RPO: 5 minutes | RTO: 30 minutes"
    for item in [
        "Active-passive failover across regions",
        "Automated database backups every 6 hours",
        "Chaos engineering tests run monthly",
        "Runbook automation via PagerDuty + Terraform",
    ]:
        p = body7.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 8: Roadmap & Next Steps ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Roadmap & Next Steps"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Q3 2025: Service mesh migration to Istio"
    for item in [
        "Q3 2025: GPU cluster for ML inference pipeline",
        "Q4 2025: Multi-region active-active deployment",
        "Q4 2025: FinOps optimization — target 20% cost reduction",
        "Q1 2026: Edge computing for latency-sensitive workloads",
    ]:
        p = body8.add_paragraph()
        p.text = item
        p.level = 1

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
