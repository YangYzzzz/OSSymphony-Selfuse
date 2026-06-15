"""
Initial Setup: Create a 6-slide Development Process presentation with empty slide 2
Task ID: impress_gf4_010
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_010'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text(slide, left, top, width, height, text, font_size=18, bold=False,
             color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a text box with formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ========== Slide 1: Title Slide ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Development Process"
    slide1.placeholders[1].text = "Strategic Overview & Methodology"

    # ========== Slide 2: "Our Development Cycle" - EMPTY content area ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    add_text(slide2, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
             "Our Development Cycle", font_size=36, bold=True,
             color=RGBColor(0x1E, 0x3A, 0x5F), alignment=PP_ALIGN.CENTER)
    # Content area intentionally empty - no shapes, no animations

    # ========== Slide 3: Team Overview ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text(slide3, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
             "Team Overview", font_size=32, bold=True,
             color=RGBColor(0x1E, 0x3A, 0x5F), alignment=PP_ALIGN.LEFT)

    team_info = [
        ("Engineering Lead", "Sarah Chen", "15 years experience in distributed systems"),
        ("Design Director", "Marcus Rivera", "Former UX lead at Figma"),
        ("QA Manager", "Priya Sharma", "Specialized in automated testing frameworks"),
        ("DevOps Lead", "James O'Brien", "AWS & Kubernetes certified architect"),
    ]
    y_pos = Inches(1.8)
    for title, name, desc in team_info:
        add_text(slide3, Inches(1.0), y_pos, Inches(3.0), Inches(0.5),
                 title, font_size=14, bold=True, color=RGBColor(0x25, 0x63, 0xEB))
        add_text(slide3, Inches(4.2), y_pos, Inches(3.0), Inches(0.5),
                 name, font_size=14, bold=False)
        add_text(slide3, Inches(7.5), y_pos, Inches(5.0), Inches(0.5),
                 desc, font_size=12, bold=False, color=RGBColor(0x66, 0x66, 0x66))
        y_pos += Inches(1.1)

    # ========== Slide 4: Project Timeline ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text(slide4, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
             "Project Timeline - Q2 2026", font_size=32, bold=True,
             color=RGBColor(0x1E, 0x3A, 0x5F), alignment=PP_ALIGN.LEFT)

    milestones = [
        ("April 7", "Sprint 1 kickoff - Requirements gathering"),
        ("April 21", "Sprint 1 review - Core architecture finalized"),
        ("May 5", "Sprint 2 kickoff - Feature development begins"),
        ("May 19", "Sprint 2 review - Alpha release candidate"),
        ("June 2", "Sprint 3 kickoff - Integration testing"),
        ("June 16", "Sprint 3 review - Beta release"),
        ("June 30", "Production deployment & monitoring"),
    ]
    y_pos = Inches(1.8)
    for date, milestone in milestones:
        add_text(slide4, Inches(1.0), y_pos, Inches(2.0), Inches(0.4),
                 date, font_size=14, bold=True, color=RGBColor(0x25, 0x63, 0xEB))
        add_text(slide4, Inches(3.5), y_pos, Inches(8.5), Inches(0.4),
                 milestone, font_size=14, bold=False)
        y_pos += Inches(0.7)

    # ========== Slide 5: Technology Stack ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text(slide5, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
             "Technology Stack", font_size=32, bold=True,
             color=RGBColor(0x1E, 0x3A, 0x5F), alignment=PP_ALIGN.LEFT)

    categories = [
        ("Frontend", "React 18, TypeScript 5.4, Tailwind CSS, Vite"),
        ("Backend", "Python 3.12, FastAPI, SQLAlchemy 2.0, Redis"),
        ("Infrastructure", "AWS EKS, Terraform, GitHub Actions, DataDog"),
        ("Database", "PostgreSQL 16, ElasticSearch 8.x, DynamoDB"),
    ]
    y_pos = Inches(2.0)
    for cat, tech in categories:
        add_text(slide5, Inches(1.0), y_pos, Inches(3.0), Inches(0.5),
                 cat, font_size=16, bold=True, color=RGBColor(0x1E, 0x3A, 0x5F))
        add_text(slide5, Inches(4.5), y_pos, Inches(8.0), Inches(0.5),
                 tech, font_size=14, bold=False, color=RGBColor(0x33, 0x33, 0x33))
        y_pos += Inches(1.1)

    # ========== Slide 6: Key Metrics ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text(slide6, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
             "Key Performance Metrics", font_size=32, bold=True,
             color=RGBColor(0x1E, 0x3A, 0x5F), alignment=PP_ALIGN.LEFT)

    metrics = [
        ("99.95%", "Uptime SLA Target"),
        ("< 200ms", "API Response Time (p95)"),
        ("85%", "Code Coverage Requirement"),
        ("< 4 hrs", "Mean Time to Recovery"),
    ]
    x_pos = Inches(0.8)
    for value, label in metrics:
        add_text(slide6, x_pos, Inches(2.5), Inches(2.8), Inches(1.0),
                 value, font_size=36, bold=True,
                 color=RGBColor(0x25, 0x63, 0xEB), alignment=PP_ALIGN.CENTER)
        add_text(slide6, x_pos, Inches(3.8), Inches(2.8), Inches(0.8),
                 label, font_size=14, bold=False,
                 color=RGBColor(0x66, 0x66, 0x66), alignment=PP_ALIGN.CENTER)
        x_pos += Inches(3.1)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
