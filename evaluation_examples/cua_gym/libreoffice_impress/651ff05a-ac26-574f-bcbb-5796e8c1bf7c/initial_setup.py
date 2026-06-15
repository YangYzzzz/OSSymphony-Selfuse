"""
Initial Setup: Create a 20-slide Manual Guide presentation with TOC on slide 2 (no hyperlinks).
Task ID: impress_fix_064
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_064'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


# Chapter structure: 6 chapters, each starting at specific slides
CHAPTERS = [
    ("Chapter 1: Introduction to the System", 3),
    ("Chapter 2: Installation and Setup", 6),
    ("Chapter 3: User Interface Overview", 9),
    ("Chapter 4: Configuration and Customization", 12),
    ("Chapter 5: Troubleshooting Common Issues", 15),
    ("Chapter 6: Advanced Features and Tips", 18),
]

# Content for each chapter's slides (3 slides per chapter)
CHAPTER_CONTENT = {
    "Chapter 1: Introduction to the System": [
        {
            "title": "Introduction to the System",
            "bullets": [
                "Welcome to the Enterprise Resource Management Platform v4.2",
                "Designed for teams of 10 to 10,000 employees",
                "Integrates with existing HR, Finance, and Operations workflows",
                "Available on Windows, macOS, and Linux",
            ],
        },
        {
            "title": "System Requirements",
            "bullets": [
                "Minimum 8 GB RAM, recommended 16 GB",
                "Processor: Intel i5 / AMD Ryzen 5 or higher",
                "Storage: 2 GB free disk space for application files",
                "Network: Stable broadband connection (10 Mbps+)",
            ],
        },
        {
            "title": "Getting Started",
            "bullets": [
                "Visit portal.erm-platform.com to create your organization account",
                "Download the desktop client from the Downloads section",
                "Contact your IT administrator for license activation",
                "Complete the initial setup wizard in under 5 minutes",
            ],
        },
    ],
    "Chapter 2: Installation and Setup": [
        {
            "title": "Installation Steps",
            "bullets": [
                "Run the installer package as administrator",
                "Select installation directory (default: C:\\Program Files\\ERM)",
                "Choose components: Core, Analytics, Reporting modules",
                "Estimated install time: 3-5 minutes on SSD",
            ],
        },
        {
            "title": "Initial Configuration",
            "bullets": [
                "Enter organization ID and admin credentials",
                "Configure database connection (PostgreSQL 14+ recommended)",
                "Set backup schedule: daily, weekly, or custom intervals",
                "Enable multi-factor authentication for all admin accounts",
            ],
        },
        {
            "title": "Environment Setup",
            "bullets": [
                "Development: localhost:8080 with debug logging enabled",
                "Staging: staging.erm-platform.com with production-like data",
                "Production: Deployed behind load balancer with SSL termination",
                "API rate limits: 1000 requests/minute per organization",
            ],
        },
    ],
    "Chapter 3: User Interface Overview": [
        {
            "title": "Dashboard Navigation",
            "bullets": [
                "Main dashboard displays KPIs, recent activity, and alerts",
                "Left sidebar: Modules, Settings, User Management",
                "Top bar: Search, Notifications, Profile settings",
                "Customizable widget layout with drag-and-drop support",
            ],
        },
        {
            "title": "Data Views and Filters",
            "bullets": [
                "Table view: Sort, filter, and export up to 50,000 records",
                "Kanban board: Drag tasks between status columns",
                "Calendar view: Visualize deadlines and milestones",
                "Saved filters: Create and share custom filter presets",
            ],
        },
        {
            "title": "Keyboard Shortcuts",
            "bullets": [
                "Ctrl+N: Create new record in current module",
                "Ctrl+F: Open global search overlay",
                "Ctrl+Shift+E: Export current view to CSV",
                "F1: Open context-sensitive help documentation",
            ],
        },
    ],
    "Chapter 4: Configuration and Customization": [
        {
            "title": "System Preferences",
            "bullets": [
                "Theme: Light, Dark, or System-default appearance",
                "Language: 24 supported languages including CJK",
                "Date format: ISO 8601, US, EU, or custom patterns",
                "Notification preferences: Email, in-app, Slack, Teams",
            ],
        },
        {
            "title": "Custom Fields and Workflows",
            "bullets": [
                "Create custom fields: text, number, date, dropdown, multi-select",
                "Define approval workflows with up to 10 sequential stages",
                "Conditional logic: Show/hide fields based on record status",
                "Field validation rules: regex, range, required, unique",
            ],
        },
        {
            "title": "Permissions and Roles",
            "bullets": [
                "Built-in roles: Admin, Manager, Editor, Viewer",
                "Custom roles with granular module-level permissions",
                "Row-level security: Restrict access by department or region",
                "Audit log: Track all permission changes with timestamps",
            ],
        },
    ],
    "Chapter 5: Troubleshooting Common Issues": [
        {
            "title": "Login and Authentication Problems",
            "bullets": [
                "Forgot password: Use self-service reset via registered email",
                "Account locked: Contact admin after 5 failed login attempts",
                "SSO errors: Verify SAML configuration in identity provider",
                "Session timeout: Default 30 minutes, configurable up to 8 hours",
            ],
        },
        {
            "title": "Performance Optimization",
            "bullets": [
                "Slow dashboard: Check browser extensions and cache size",
                "Large dataset lag: Enable server-side pagination in settings",
                "Report generation timeout: Reduce date range or add filters",
                "Database connection drops: Verify network stability and pool size",
            ],
        },
        {
            "title": "Data Import/Export Issues",
            "bullets": [
                "CSV encoding: Use UTF-8 with BOM for international characters",
                "Column mapping: Match headers exactly (case-sensitive)",
                "Maximum file size: 50 MB per upload, split larger files",
                "Duplicate detection: Enable unique key validation before import",
            ],
        },
    ],
    "Chapter 6: Advanced Features and Tips": [
        {
            "title": "API Integration",
            "bullets": [
                "RESTful API with OpenAPI 3.0 specification",
                "Authentication: OAuth 2.0 with bearer tokens",
                "Webhooks: Configure event-driven notifications",
                "Rate limiting: 1000 req/min standard, 5000 req/min enterprise",
            ],
        },
        {
            "title": "Automation and Scripting",
            "bullets": [
                "Built-in macro language for repetitive task automation",
                "Scheduled jobs: Run reports and data syncs on cron schedules",
                "Email templates: Dynamic field insertion with Jinja2 syntax",
                "Workflow triggers: Execute actions on record create/update/delete",
            ],
        },
        {
            "title": "Best Practices and Tips",
            "bullets": [
                "Regular backups: Test restore procedure quarterly",
                "Version control: Tag configuration changes with release notes",
                "User training: Schedule onboarding sessions for new team members",
                "Stay updated: Subscribe to changelog at docs.erm-platform.com",
            ],
        },
    ],
}


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    # Title text
    txBox = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Enterprise Resource Management Platform"
    run.font.name = "Calibri"
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Subtitle
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "User Manual & Configuration Guide — Version 4.2"
    run2.font.name = "Calibri"
    run2.font.size = Pt(22)
    run2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)

    # --- Slide 2: Table of Contents (NO hyperlinks) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])

    # TOC heading
    txTitle = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1))
    tf_title = txTitle.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.alignment = PP_ALIGN.CENTER
    run_t = p_title.add_run()
    run_t.text = "Table of Contents"
    run_t.font.name = "Calibri"
    run_t.font.size = Pt(32)
    run_t.font.bold = True
    run_t.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    # Chapter list items (each in its own text box for clarity)
    y_start = Inches(1.8)
    for i, (chapter_title, slide_num) in enumerate(CHAPTERS):
        txBox = slide2.shapes.add_textbox(
            Inches(2), y_start + Inches(i * 0.8), Inches(9), Inches(0.7)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{chapter_title}"
        run.font.name = "Calibri"
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0x2E, 0x56, 0x84)

        # Add slide reference text
        run2 = p.add_run()
        run2.text = f"  (Slide {slide_num})"
        run2.font.name = "Calibri"
        run2.font.size = Pt(16)
        run2.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

    # --- Slides 3 through 20: Chapter content ---
    for chapter_title, start_slide in CHAPTERS:
        slides_content = CHAPTER_CONTENT[chapter_title]
        for j, content in enumerate(slides_content):
            slide = prs.slides.add_slide(prs.slide_layouts[5])

            # Chapter section header background for first slide of each chapter
            if j == 0:
                bg_fill = slide.background.fill
                bg_fill.solid()
                bg_fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF4)

            # Title
            txTitle = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(1))
            tf_t = txTitle.text_frame
            tf_t.word_wrap = True
            p_t = tf_t.paragraphs[0]
            p_t.alignment = PP_ALIGN.LEFT
            run_t = p_t.add_run()
            run_t.text = content["title"]
            run_t.font.name = "Calibri"
            run_t.font.size = Pt(28)
            run_t.font.bold = True
            run_t.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

            # Bullet content
            txBody = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.933), Inches(4.5))
            tf_b = txBody.text_frame
            tf_b.word_wrap = True
            for k, bullet in enumerate(content["bullets"]):
                if k == 0:
                    p_b = tf_b.paragraphs[0]
                else:
                    p_b = tf_b.add_paragraph()
                p_b.space_after = Pt(12)
                run_b = p_b.add_run()
                run_b.text = f"• {bullet}"
                run_b.font.name = "Calibri"
                run_b.font.size = Pt(18)
                run_b.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

            # Slide number at bottom right
            txNum = slide.shapes.add_textbox(
                Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.5)
            )
            tf_n = txNum.text_frame
            p_n = tf_n.paragraphs[0]
            p_n.alignment = PP_ALIGN.RIGHT
            run_n = p_n.add_run()
            # Slide number: 1-indexed (slide1=title, slide2=TOC, so content starts at 3)
            run_n.text = str(start_slide + j)
            run_n.font.name = "Calibri"
            run_n.font.size = Pt(12)
            run_n.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    # Verify we have exactly 20 slides
    assert len(prs.slides) == 20, f"Expected 20 slides, got {len(prs.slides)}"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
