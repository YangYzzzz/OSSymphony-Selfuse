"""
Reward Script: Verify hyperlinks on slide 2 TOC in Manual_Guide.pptx
Task ID: impress_fix_064
Domain: libreoffice_impress
Scoring: 6 components (one per chapter hyperlink), ~0.167 points each = 1.0 total
"""

import os
import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_064'

# Expected hyperlink mappings: chapter index (0-based among the 6 chapter shapes) -> target slide XML
# Chapter 1 -> slide 3, Chapter 2 -> slide 6, ... Chapter 6 -> slide 18
EXPECTED_LINKS = {
    0: 'slide3.xml',
    1: 'slide6.xml',
    2: 'slide9.xml',
    3: 'slide12.xml',
    4: 'slide15.xml',
    5: 'slide18.xml',
}

CHAPTER_LABELS = [
    'Chapter 1',
    'Chapter 2',
    'Chapter 3',
    'Chapter 4',
    'Chapter 5',
    'Chapter 6',
]

POINTS_PER_LINK = 1.0 / 6.0  # ~0.1667


def get_slide2_rels(pptx_path):
    """Parse slide2.xml.rels to build rId -> Target mapping."""
    rels = {}
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            with zf.open('ppt/slides/_rels/slide2.xml.rels') as f:
                root = ET.parse(f).getroot()
                ns = 'http://schemas.openxmlformats.org/package/2006/relationships'
                for rel in root.findall(f'{{{ns}}}Relationship'):
                    rid = rel.get('Id')
                    target = rel.get('Target')
                    rel_type = rel.get('Type', '')
                    # Only include slide-type relationships (not slideLayout)
                    if 'relationships/slide' in rel_type and 'slideLayout' not in rel_type:
                        rels[rid] = target
    except Exception as e:
        print(f"ERROR: Could not parse slide2 rels: {e}")
    return rels


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: at least 20 slides
    if len(prs.slides) < 20:
        print(f"CRITICAL: Expected at least 20 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Get slide 2 (index 1)
    slide = prs.slides[1]

    # Build rId -> target slide mapping from rels
    rels_map = get_slide2_rels(file_path)
    print(f"INFO: Slide 2 rels map: {rels_map}")

    # Collect chapter text shapes (shapes 2-7 = indices 2..7, which are the 6 chapter text boxes)
    # We identify them by text containing "Chapter N"
    chapter_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            full_text = shape.text_frame.text
            for label in CHAPTER_LABELS:
                if label in full_text:
                    chapter_shapes.append((label, shape))
                    break

    print(f"INFO: Found {len(chapter_shapes)} chapter shapes on slide 2")

    if len(chapter_shapes) < 6:
        print(f"WARNING: Expected 6 chapter shapes, found {len(chapter_shapes)}")

    # Verify each chapter hyperlink
    for idx, (label, shape) in enumerate(chapter_shapes):
        expected_target = EXPECTED_LINKS.get(idx)
        if expected_target is None:
            continue

        # Component N: Chapter N hyperlink points to correct slide
        try:
            # Find the first run with meaningful text (the chapter title run)
            hyperlink_found = False
            correct_target = False

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not (run.text or "").strip():
                        continue
                    if label not in run.text:
                        continue

                    # Check for hlinkClick element in the run's rPr
                    rPr = run._r.find(qn('a:rPr'))
                    if rPr is not None:
                        hlinkClick = rPr.find(qn('a:hlinkClick'))
                        if hlinkClick is not None:
                            action = hlinkClick.get('action', '')
                            r_id_key = qn('r:id')
                            r_id = hlinkClick.get(r_id_key, '')

                            if action == 'ppaction://hlinksldjump':
                                hyperlink_found = True
                                # Verify the relationship target
                                actual_target = rels_map.get(r_id, 'UNKNOWN')
                                if actual_target == expected_target:
                                    correct_target = True

            if hyperlink_found and correct_target:
                print(f"PASS: {label} — hyperlink to {expected_target} ({POINTS_PER_LINK:.4f} pts)")
                total_score += POINTS_PER_LINK
            elif hyperlink_found:
                actual = rels_map.get(r_id, 'UNKNOWN') if 'r_id' in dir() else 'UNKNOWN'
                print(f"FAIL: {label} — hyperlink found but targets wrong slide (expected {expected_target})")
            else:
                print(f"FAIL: {label} — no internal hyperlink (ppaction://hlinksldjump) found")

        except Exception as e:
            print(f"ERROR: {label} — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.4f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
