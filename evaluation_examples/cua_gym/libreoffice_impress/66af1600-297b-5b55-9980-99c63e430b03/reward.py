"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 8 of my LibreOffice Impress deck, I want to insert a 3×3 table and set the whole first row to bold for the headers. How do I get that done?
Generated: 2025-09-10 12:01:02
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation

# ---------------- Verification Helpers ---------------- #

def check_table_size(table, rows_required=3, cols_required=3):
    """Return True if the table has at least the required rows and cols."""
    rows = len(table.rows)
    cols = len(table.columns)
    print(f"    Table size found: {rows} rows x {cols} cols")
    return rows >= rows_required and cols >= cols_required


def check_first_row_bold(table):
    """Return True if every run of text in the first row is bold."""
    first_row = table.rows[0]
    all_cells_bold = True

    for idx, cell in enumerate(first_row.cells):
        cell_text = cell.text_frame.text.strip()
        if not cell_text:
            # Empty header cell = failure for bold requirement
            print(f"      Cell {idx} is empty – header missing (fail)")
            all_cells_bold = False
            continue

        cell_has_bold = False
        # Inspect every run for bold formatting
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                # Log run details
                if run.text.strip():
                    print(f"        Run '{run.text}' bold={run.font.bold}")
                # A run counts as bold if run.font.bold is explicitly True
                if run.text.strip() and run.font.bold:
                    cell_has_bold = True
        if not cell_has_bold:
            print(f"      Cell {idx} does not contain bold text")
            all_cells_bold = False
    return all_cells_bold

# ---------------- Main Verification Function ---------------- #

def verify_impress_table_task(file_path):
    """
    Verify that on slide 8 of the presentation there is a 3×3 table
    and that the entire first row is bold. Progressive scoring:
        • 0.4 pts – Required-size table exists on slide 8
        • 0.6 pts – All first-row header cells are bold
    Returns a float score between 0.0 and 1.0.
    """
    print(f"Checking presentation at: {file_path}")

    total_score = 0.0
    max_score = 1.0

    # ---------- Basic file checks (no points awarded) ---------- #
    if not os.path.exists(file_path):
        print("✗ File does not exist")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # ---------- Slide 8 existence ---------- #
    if len(prs.slides) < 8:
        print(f"✗ Presentation has only {len(prs.slides)} slides; slide 8 missing")
        print("REWARD: 0.0")
        return 0.0
    else:
        print(f"✓ Presentation has {len(prs.slides)} slides (slide 8 exists)")

    slide8 = prs.slides[7]  # zero-based index

    # ---------- Table existence & size ---------- #
    matching_table = None
    for shape in slide8.shapes:
        if shape.has_table:
            tbl = shape.table
            if check_table_size(tbl):
                matching_table = tbl
                break  # first qualifying table is enough

    if matching_table:
        print("✓ Found table with required size on slide 8 (0.4 points)")
        total_score += 0.4
    else:
        print("✗ No 3×3 (or larger) table found on slide 8")
        print(f"REWARD: {total_score}")
        return total_score  # cannot continue bold check without table

    # ---------- Bold header verification ---------- #
    if check_first_row_bold(matching_table):
        print("✓ First-row headers are bold (0.6 points)")
        total_score += 0.6
    else:
        print("✗ First-row headers are not uniformly bold")

    # ---------- Final score ---------- #
    final_score = min(total_score, max_score)
    print(f"Total score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score

# ---------------- Execute when run directly ---------------- #
if __name__ == "__main__":
    FILE_PATH = "/home/user/on_slide_8_of_my_libreoffice_impress_deck_i_want_to_insert_a_33_table_and_set_the_whole_first_row_to_golden.pptx"
    verify_impress_table_task(FILE_PATH)
