"""
FINAL REWARD SCRIPT - SUCCESS
Task: The final slide in my LibreOffice Impress deck ended up with a jumble of fonts after some copy-pasting. How can I quickly grab every text box on that one slide (slide 12) and force them all to use Times New Roman so everything matches?
Generated: 2025-09-10 12:36:10
Status: success
Model: azure-o3
Total Steps: 4
"""

from pptx import Presentation
import os

def verify_times_new_roman_slide(file_path: str, slide_number: int = 12) -> float:
    """Verify that every text run on the specified slide uses Times New Roman.

    Progressive scoring:
        • 1.0  – All text runs use Times New Roman
        • 0.8  – ≥80 % of text runs use Times New Roman
        • 0.6  – ≥50 % of text runs use Times New Roman
        • 0.2  – >0 % (but <50 %) use Times New Roman
        • 0.0  – No text runs or none use Times New Roman
    """

    print(f"Verifying slide {slide_number} fonts in: {file_path}")

    # ---------- 0.  File existence & loading (no points for this – prerequisite) ----------
    if not os.path.exists(file_path):
        print("✗ Presentation file not found")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        return 0.0

    # ---------- 1.  Slide existence check (still prerequisite – no points) ----------
    if slide_number < 1 or slide_number > len(prs.slides):
        print(f"✗ Slide {slide_number} does not exist (presentation has {len(prs.slides)} slides)")
        return 0.0

    slide = prs.slides[slide_number - 1]

    # ---------- 2.  Inspect every text run on the slide ----------
    total_runs = 0
    matching_runs = 0
    other_fonts = set()

    for shape in slide.shapes:
        if not hasattr(shape, "has_text_frame") or not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text or not run.text.strip():
                    continue  # ignore empty runs
                total_runs += 1

                # Primary font name recorded by python-pptx
                font_name = run.font.name

                # Fallback: sometimes the rPr attribute carries the typeface
                if font_name is None:
                    rPr = run._r.get_or_add_rPr()
                    font_name = rPr.get("typeface")  # may still be None

                if font_name and "times" in font_name.lower() and "new" in font_name.lower():
                    matching_runs += 1
                else:
                    other_fonts.add(font_name or "UNKNOWN")

    # ---------- 3.  Scoring ----------
    if total_runs == 0:
        print("✗ No text found on slide – cannot verify fonts")
        return 0.0

    ratio = matching_runs / total_runs
    print(f"Total text runs: {total_runs}")
    print(f"Runs using Times New Roman: {matching_runs}")
    print(f"Matching ratio: {ratio:.2%}")
    if other_fonts:
        print(f"Non-Times New Roman fonts detected: {other_fonts}")

    # Progressive score based on ratio
    if ratio == 1.0:
        score = 1.0
    elif ratio >= 0.8:
        score = 0.8
    elif ratio >= 0.5:
        score = 0.6
    elif ratio > 0:
        score = 0.2
    else:
        score = 0.0

    print(f"Calculated score: {score}")
    return score

# --------------- Run verification when executed as a script ---------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/the_final_slide_in_my_libreoffice_impress_deck_ended_up_with_a_jumble_of_fonts_after_some_copy_pasti_golden.pptx"
    reward = verify_times_new_roman_slide(FILE_PATH)
    print(f"REWARD: {reward}")

