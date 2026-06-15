"""
Initial Setup: Create a portfolio presentation with 15 slides and one master slide named 'Default'.
Task ID: impress_ma_005
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_005'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Rename the default slide master to 'Default'
    master = prs.slide_masters[0]
    # The master name is stored in the slideMaster XML as the 'name' attribute on cSld
    # But in LibreOffice, master slide names come from the slide master XML
    # python-pptx doesn't expose master name directly, so we set it via XML
    cSld = master._element.find(qn('p:cSld'))
    if cSld is not None:
        cSld.set('name', 'Default')

    # Color palette for the portfolio
    DARK_BG = RGBColor(0x2C, 0x3E, 0x50)
    ACCENT = RGBColor(0xE7, 0x4C, 0x3C)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
    DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)

    # --- Slide 1: Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(2),
                 "Elena Rodriguez", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
                 "Senior UX Designer & Creative Director", font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5), Inches(11), Inches(1),
                 "elena.rodriguez@designstudio.com | +1 (415) 555-0142", font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

    # --- Slide 2: About Me ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1),
                 "About Me", font_size=36, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.LEFT)
    about_text = (
        "With over 12 years of experience in UX design and creative direction, "
        "I specialize in building intuitive digital experiences for enterprise "
        "and consumer products. My work spans fintech, healthcare, and e-commerce "
        "platforms, serving millions of users worldwide."
    )
    add_text_box(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(3),
                 about_text, font_size=18, color=DARK_TEXT)

    # --- Slide 3: Skills Overview ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1),
                 "Core Skills", font_size=36, bold=True, color=DARK_TEXT)
    skills = [
        "User Research & Persona Development",
        "Interaction Design & Prototyping",
        "Design Systems & Component Libraries",
        "Accessibility & Inclusive Design (WCAG 2.1)",
        "Cross-Platform Design (iOS, Android, Web)",
        "Stakeholder Management & Design Sprints",
    ]
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, skill in enumerate(skills):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  \u2022  {skill}"
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.size = Pt(18)
        run.font.color.rgb = DARK_TEXT

    # --- Slide 4: Project 1 - FinPay Dashboard ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1),
                 "Project: FinPay Dashboard Redesign", font_size=32, bold=True, color=DARK_TEXT)
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(5), Inches(0.6),
                 "Client: FinPay Technologies | 2024", font_size=14, color=ACCENT)
    add_text_box(slide, Inches(0.5), Inches(2.3), Inches(12), Inches(3),
                 "Redesigned the payment dashboard for 2.3M active users. Reduced task completion "
                 "time by 34% and increased user satisfaction scores from 3.2 to 4.7 out of 5. "
                 "Led a team of 4 designers through a 12-week sprint cycle.",
                 font_size=16, color=DARK_TEXT)

    # --- Slide 5: Project 1 - Metrics ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1),
                 "FinPay: Key Results", font_size=32, bold=True, color=DARK_TEXT)
    metrics = [
        ("Task Completion Rate", "+34%", "Streamlined 5-step flow to 3 steps"),
        ("User Satisfaction", "4.7/5.0", "Up from 3.2 baseline score"),
        ("Support Tickets", "-52%", "Clearer UI reduced confusion"),
        ("Onboarding Time", "2.1 min", "Down from 8.5 minutes average"),
    ]
    table_shape = slide.shapes.add_table(len(metrics) + 1, 3, Inches(0.5), Inches(1.8), Inches(12), Inches(3))
    table = table_shape.table
    for i, header in enumerate(["Metric", "Result", "Detail"]):
        cell = table.cell(0, i)
        cell.text = header
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, (metric, result, detail) in enumerate(metrics, 1):
        table.cell(r, 0).text = metric
        table.cell(r, 1).text = result
        table.cell(r, 2).text = detail

    # --- Slide 6: Project 2 - MedConnect ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1),
                 "Project: MedConnect Patient Portal", font_size=32, bold=True, color=DARK_TEXT)
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(5), Inches(0.6),
                 "Client: Valley Health Network | 2023-2024", font_size=14, color=ACCENT)
    add_text_box(slide, Inches(0.5), Inches(2.3), Inches(12), Inches(3),
                 "Designed an accessible patient portal compliant with WCAG 2.1 AA standards. "
                 "Features include appointment scheduling, lab results viewing, secure messaging, "
                 "and prescription management. Served 450K patients across 12 clinic locations.",
                 font_size=16, color=DARK_TEXT)

    # --- Slide 7: Project 2 - Design Process ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1),
                 "MedConnect: Design Process", font_size=32, bold=True, color=DARK_TEXT)
    steps = [
        "1. Discovery: 45 patient interviews, 12 clinician shadowing sessions",
        "2. Define: Mapped 8 key user journeys, identified 23 pain points",
        "3. Ideate: Facilitated 6 design thinking workshops with stakeholders",
        "4. Prototype: Built high-fidelity Figma prototype with 120+ screens",
        "5. Test: 3 rounds of usability testing with 60 participants",
        "6. Launch: Phased rollout across 12 locations over 4 months",
    ]
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, step in enumerate(steps):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = step
        p.space_after = Pt(10)
        run = p.runs[0]
        run.font.size = Pt(16)
        run.font.color.rgb = DARK_TEXT

    # --- Slide 8: Project 3 - ShopEase ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1),
                 "Project: ShopEase Mobile App", font_size=32, bold=True, color=DARK_TEXT)
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(5), Inches(0.6),
                 "Client: ShopEase Inc. | 2023", font_size=14, color=ACCENT)
    add_text_box(slide, Inches(0.5), Inches(2.3), Inches(12), Inches(3),
                 "Designed a native mobile shopping experience for iOS and Android. "
                 "Implemented AR try-on features for fashion items. App achieved 4.8-star "
                 "rating with 1.2M downloads in the first quarter after launch.",
                 font_size=16, color=DARK_TEXT)

    # --- Slide 9: Design Philosophy ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1),
                 "Design Philosophy", font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    quotes = [
        "\u201cGreat design is invisible \u2014 it feels effortless to the user.\u201d",
        "\u201cAccessibility is not an afterthought; it's the foundation.\u201d",
        "\u201cData informs design, but empathy drives it.\u201d",
    ]
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, quote in enumerate(quotes):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p = tf.add_paragraph()  # extra spacing
        p.text = quote
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(22)
        run.font.italic = True
        run.font.color.rgb = LIGHT_GRAY

    # --- Slide 10: Tools & Technologies ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1),
                 "Tools & Technologies", font_size=36, bold=True, color=DARK_TEXT)
    tools_data = [
        ("Design", "Figma, Sketch, Adobe XD, InVision"),
        ("Prototyping", "Figma, Principle, ProtoPie, Framer"),
        ("Research", "UserTesting, Hotjar, Maze, Dovetail"),
        ("Collaboration", "Miro, FigJam, Notion, Confluence"),
        ("Development", "HTML/CSS, React basics, Storybook"),
    ]
    table_shape = slide.shapes.add_table(len(tools_data) + 1, 2, Inches(0.5), Inches(1.8), Inches(12), Inches(3.5))
    table = table_shape.table
    table.columns[0].width = Inches(3)
    table.columns[1].width = Inches(9)
    table.cell(0, 0).text = "Category"
    table.cell(0, 1).text = "Tools"
    for run in table.cell(0, 0).text_frame.paragraphs[0].runs:
        run.font.bold = True
    for run in table.cell(0, 1).text_frame.paragraphs[0].runs:
        run.font.bold = True
    for r, (cat, tools) in enumerate(tools_data, 1):
        table.cell(r, 0).text = cat
        table.cell(r, 1).text = tools

    # --- Slide 11: Education ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1),
                 "Education & Certifications", font_size=36, bold=True, color=DARK_TEXT)
    edu_items = [
        "MFA in Interaction Design \u2014 School of Visual Arts, New York (2013)",
        "BA in Graphic Design \u2014 University of California, Berkeley (2011)",
        "Google UX Design Professional Certificate (2022)",
        "Certified Usability Analyst (CUA) \u2014 Human Factors International",
        "IAAP Web Accessibility Specialist (WAS) Certification",
    ]
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(edu_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  \u2022  {item}"
        p.space_after = Pt(10)
        run = p.runs[0]
        run.font.size = Pt(16)
        run.font.color.rgb = DARK_TEXT

    # --- Slide 12: Testimonials ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1),
                 "Client Testimonials", font_size=36, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
                 "\u201cElena transformed our product experience. Her attention to detail "
                 "and ability to translate complex requirements into elegant solutions "
                 "is unmatched.\u201d\n\u2014 James Liu, VP Product, FinPay Technologies",
                 font_size=16, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4), Inches(11), Inches(1.5),
                 "\u201cWorking with Elena was a game-changer for our patient portal. "
                 "She brought empathy and rigor to every design decision.\u201d\n"
                 "\u2014 Dr. Sarah Kim, Chief Digital Officer, Valley Health Network",
                 font_size=16, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

    # --- Slide 13: Awards ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1),
                 "Awards & Recognition", font_size=36, bold=True, color=DARK_TEXT)
    awards = [
        "Red Dot Design Award \u2014 Best UX, FinPay Dashboard (2024)",
        "Webby Award \u2014 Best Health App, MedConnect Portal (2024)",
        "IXDA Interaction Award \u2014 Honorable Mention (2023)",
        "Fast Company Innovation by Design \u2014 Finalist (2023)",
        "A\u2019 Design Award \u2014 Gold, ShopEase Mobile App (2023)",
    ]
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, award in enumerate(awards):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  \u2022  {award}"
        p.space_after = Pt(10)
        run = p.runs[0]
        run.font.size = Pt(16)
        run.font.color.rgb = DARK_TEXT

    # --- Slide 14: Speaking & Publications ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1),
                 "Speaking & Publications", font_size=36, bold=True, color=DARK_TEXT)
    speaking = [
        "Keynote: \u201cDesigning for Trust in FinTech\u201d \u2014 UX London 2024",
        "Workshop: \u201cAccessible Design Systems\u201d \u2014 Config 2024",
        "Panel: \u201cAI in UX Research\u201d \u2014 SXSW Interactive 2023",
        "Article: \u201c5 Principles for Healthcare UX\u201d \u2014 UX Magazine",
        "Co-author: \u201cDesign Systems at Scale\u201d \u2014 A Book Apart (2023)",
    ]
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(speaking):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  \u2022  {item}"
        p.space_after = Pt(10)
        run = p.runs[0]
        run.font.size = Pt(16)
        run.font.color.rgb = DARK_TEXT

    # --- Slide 15: Contact ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(1),
                 "Let\u2019s Connect", font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    contact_info = (
        "Email: elena.rodriguez@designstudio.com\n"
        "Phone: +1 (415) 555-0142\n"
        "Portfolio: www.elenadesigns.com\n"
        "LinkedIn: linkedin.com/in/elenarodriguez"
    )
    add_text_box(slide, Inches(2), Inches(3), Inches(9), Inches(3),
                 contact_info, font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
