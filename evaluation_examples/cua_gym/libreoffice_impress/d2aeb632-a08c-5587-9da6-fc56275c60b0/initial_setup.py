"""
Initial Setup: Marketing_Slides.pptx with 8 slides, various transitions, and custom font usage.
Task ID: osworld_multi_apps_impress_infeasible_006
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import zipfile
import copy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'osworld_multi_apps_impress_infeasible_006'
OUTPUT = f'{WORKDIR}/Marketing_Slides.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_transition_to_slide(prs, slide_idx, transition_type):
    """Add a transition effect to a slide via XML manipulation."""
    # Get the pptx as a zip and modify the slide XML
    # We do this by modifying the slide's XML element directly
    slide = prs.slides[slide_idx]
    spTree = slide.shapes._spTree
    cSld = spTree.getparent()

    # Build transition XML
    nsmap = 'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
    transitions = {
        'fade':     '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>',
        'dissolve': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:dissolve/></p:transition>',
        'push':     '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med" dir="l"><p:push/></p:transition>',
        'wipe':     '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:wipe/></p:transition>',
        'zoom':     '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:zoom/></p:transition>',
        'blinds':   '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:blinds dir="horz"/></p:transition>',
        'checker':  '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:checker dir="vert"/></p:transition>',
        'cover':    '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med" dir="l"><p:cover/></p:transition>',
    }
    xml_str = transitions.get(transition_type, transitions['fade'])

    ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ET.register_namespace('p', ns)

    tr_elem = ET.fromstring(xml_str)
    cSld.append(tr_elem)


def create_initial():
    prs = Presentation()

    # Slide dimensions: standard widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_layouts = prs.slide_layouts

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(slide_layouts[0])  # Title Slide layout
    slide1.shapes.title.text = "Q1 2025 Marketing Strategy"
    slide1.placeholders[1].text = "Accelerating Growth Through Digital Channels"

    title_tf = slide1.shapes.title.text_frame
    for para in title_tf.paragraphs:
        for run in para.runs:
            run.font.name = "Georgia"
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

    sub_tf = slide1.placeholders[1].text_frame
    for para in sub_tf.paragraphs:
        for run in para.runs:
            run.font.name = "Georgia"
            run.font.size = Pt(24)
            run.font.color.rgb = RGBColor(0x40, 0x40, 0x40)

    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = RGBColor(0xF2, 0xF7, 0xFF)

    # ---- Slide 2: Agenda ----
    slide2 = prs.slides.add_slide(slide_layouts[1])  # Title and Content
    slide2.shapes.title.text = "Today's Agenda"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Market Overview & Trends"
    items2 = [
        "Campaign Performance Review",
        "Digital Channel Analysis",
        "Competitive Landscape",
        "Q2 Strategy & Initiatives",
        "Budget Allocation",
        "Key Milestones & Timeline",
    ]
    for item in items2:
        p = tf2.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(20)

    slide2.shapes.title.text_frame.paragraphs[0].runs[0].font.name = "Georgia"
    slide2.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    slide2.shapes.title.text_frame.paragraphs[0].runs[0].font.bold = True

    # ---- Slide 3: Market Overview ----
    slide3 = prs.slides.add_slide(slide_layouts[1])
    slide3.shapes.title.text = "Market Overview"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Global Digital Ad Spend: $740B in 2024 (+12% YoY)"
    metrics = [
        "Mobile accounts for 68% of all digital traffic",
        "Video content drives 3x higher engagement rates",
        "Social commerce growing at 28% annually",
        "AI-driven personalization lifts conversion by 40%",
        "Email ROI remains highest at $42 per $1 spent",
    ]
    for metric in metrics:
        p = tf3.add_paragraph()
        p.text = metric
        p.level = 1
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(18)

    slide3.background.fill.solid()
    slide3.background.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)

    # ---- Slide 4: Campaign Performance ----
    slide4 = prs.slides.add_slide(slide_layouts[5])  # Blank

    # Title text box
    txbox_title = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf_title = txbox_title.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "Q1 Campaign Performance"
    p_title.runs[0].font.name = "Georgia"
    p_title.runs[0].font.size = Pt(34)
    p_title.runs[0].font.bold = True
    p_title.runs[0].font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

    # Performance table
    table_shape = slide4.shapes.add_table(6, 4, Inches(0.5), Inches(1.5), Inches(12), Inches(4))
    table = table_shape.table
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(3.0)
    table.columns[2].width = Inches(2.5)
    table.columns[3].width = Inches(3.0)

    headers = ["Campaign", "Budget Spent", "Impressions", "Conversion Rate"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.name = "Calibri"
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # Header background
        from pptx.oxml.ns import qn
        from lxml import etree
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', '1F3864')

    campaign_data = [
        ["Summer Launch 2024", "$125,000", "8.4M", "3.2%"],
        ["Back to School", "$89,500", "6.1M", "4.1%"],
        ["Holiday Season Blast", "$210,000", "15.7M", "5.8%"],
        ["New Year Promo", "$67,200", "4.3M", "2.9%"],
        ["Valentine's Day", "$45,800", "3.2M", "3.7%"],
    ]

    for row_idx, row_data in enumerate(campaign_data, 1):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = "Calibri"
                    run.font.size = Pt(15)

    # ---- Slide 5: Digital Channels ----
    slide5 = prs.slides.add_slide(slide_layouts[1])
    slide5.shapes.title.text = "Digital Channel Analysis"
    tf5 = slide5.placeholders[1].text_frame

    channels = [
        ("Social Media (Meta, TikTok, LinkedIn)", "38% of total budget — highest reach"),
        ("Search Engine Marketing (Google, Bing)", "28% of budget — strongest intent targeting"),
        ("Programmatic Display", "18% of budget — brand awareness"),
        ("Email Marketing", "8% of budget — highest ROI channel"),
        ("Influencer Partnerships", "8% of budget — trust & authenticity"),
    ]

    tf5.text = channels[0][0]
    for run in tf5.paragraphs[0].runs:
        run.font.name = "Calibri"
        run.font.bold = True
        run.font.size = Pt(18)

    for ch_name, ch_desc in channels[1:]:
        p = tf5.add_paragraph()
        p.text = ch_name
        p.level = 0
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.bold = True
            run.font.size = Pt(18)
        p2 = tf5.add_paragraph()
        p2.text = ch_desc
        p2.level = 1
        for run in p2.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(16)
            run.font.italic = True

    slide5.background.fill.solid()
    slide5.background.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)

    # ---- Slide 6: Competitive Landscape ----
    slide6 = prs.slides.add_slide(slide_layouts[5])  # Blank

    txbox_title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    p_t6 = txbox_title6.text_frame.paragraphs[0]
    p_t6.text = "Competitive Landscape"
    p_t6.runs[0].font.name = "Georgia"
    p_t6.runs[0].font.size = Pt(34)
    p_t6.runs[0].font.bold = True
    p_t6.runs[0].font.color.rgb = RGBColor(0x7B, 0x00, 0x2A)

    # Competitor comparison text box
    txbox6 = slide6.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5))
    tf6 = txbox6.text_frame
    tf6.word_wrap = True

    competitors = [
        ("BrandCore Inc.", "Market leader, 34% share", "Strong in SEM, weak in social"),
        ("MediaFusion Ltd.", "Challenger, 18% share", "Aggressive influencer spend"),
        ("PixelMark Group", "Emerging, 11% share", "Video-first strategy"),
    ]

    for comp_name, share, strategy in competitors:
        p = tf6.add_paragraph()
        p.text = comp_name
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.bold = True
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x7B, 0x00, 0x2A)
        p2 = tf6.add_paragraph()
        p2.text = share
        p2.level = 1
        for run in p2.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(16)
        p3 = tf6.add_paragraph()
        p3.text = strategy
        p3.level = 1
        for run in p3.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(15)
            run.font.italic = True
        p4 = tf6.add_paragraph()
        p4.text = ""

    # Our position text box
    txbox6b = slide6.shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(6), Inches(5))
    tf6b = txbox6b.text_frame
    tf6b.word_wrap = True
    pb = tf6b.paragraphs[0]
    pb.text = "Our Competitive Advantages"
    for run in pb.runs:
        run.font.name = "Georgia"
        run.font.bold = True
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

    advantages = [
        "Industry-leading data analytics platform",
        "Proprietary AI personalization engine",
        "Cross-channel attribution modeling",
        "Award-winning creative production",
        "Strategic publisher partnerships",
    ]
    for adv in advantages:
        p = tf6b.add_paragraph()
        p.text = adv
        p.level = 1
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(16)

    slide6.background.fill.solid()
    slide6.background.fill.fore_color.rgb = RGBColor(0xFF, 0xF5, 0xF7)

    # ---- Slide 7: Q2 Strategy ----
    slide7 = prs.slides.add_slide(slide_layouts[1])
    slide7.shapes.title.text = "Q2 2025 Strategy & Key Initiatives"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "Deepen AI-Driven Personalization"

    initiatives = [
        "Launch predictive audience segmentation across all channels",
        "Expand video content budget by 25% — target 200M views",
        "Pilot shoppable video on TikTok and Instagram",
        "Introduce loyalty-tier email automation sequences",
        "Deploy real-time competitive spend monitoring",
        "Scale influencer micro-network from 120 to 350 creators",
    ]
    for init in initiatives:
        p = tf7.add_paragraph()
        p.text = init
        p.level = 1
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(17)

    for run in tf7.paragraphs[0].runs:
        run.font.name = "Calibri"
        run.font.bold = True
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

    slide7.background.fill.solid()
    slide7.background.fill.fore_color.rgb = RGBColor(0xF0, 0xF8, 0xF0)

    # ---- Slide 8: Closing / Call to Action ----
    slide8 = prs.slides.add_slide(slide_layouts[0])  # Title Slide layout
    slide8.shapes.title.text = "Let's Build the Future Together"
    slide8.placeholders[1].text = "Questions? Contact: marketing@company.com\nFollow-up meeting: March 15, 2025 | 10:00 AM"

    for para in slide8.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Georgia"
            run.font.size = Pt(36)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    for para in slide8.placeholders[1].text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(22)
            run.font.color.rgb = RGBColor(0xE0, 0xE8, 0xFF)

    slide8.background.fill.solid()
    slide8.background.fill.fore_color.rgb = RGBColor(0x1F, 0x38, 0x64)

    # Add speaker notes to some slides
    slide1.notes_slide.notes_text_frame.text = "Welcome audience, emphasize Q1 wins before diving into strategy."
    slide3.notes_slide.notes_text_frame.text = "Reference Gartner 2024 report for market size data."
    slide5.notes_slide.notes_text_frame.text = "Highlight TikTok as fastest growing channel for our demographic."
    slide7.notes_slide.notes_text_frame.text = "Q2 budget approval pending CFO sign-off on March 10."

    # Save presentation
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Now add transitions via XML post-processing
    add_transitions_via_xml(OUTPUT)

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


def add_transitions_via_xml(pptx_path):
    """Add slide transitions by modifying the pptx XML directly after saving."""
    import shutil
    import tempfile

    # Transition types to assign to each slide (8 slides)
    transitions_xml = [
        # Slide 1: fade
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>',
        # Slide 2: dissolve (push left as dissolve equivalent)
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade advTm="0"/></p:transition>',
        # Slide 3: wipe
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:wipe dir="l"/></p:transition>',
        # Slide 4: cover
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med" dir="l"><p:cover/></p:transition>',
        # Slide 5: zoom
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:zoom/></p:transition>',
        # Slide 6: blinds
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:blinds dir="horz"/></p:transition>',
        # Slide 7: checker
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:checker dir="vert"/></p:transition>',
        # Slide 8: fade
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="slow"><p:fade/></p:transition>',
    ]

    # Work with the zip file
    tmp_path = pptx_path + '.tmp'
    shutil.copy(pptx_path, tmp_path)

    with zipfile.ZipFile(tmp_path, 'r') as zin:
        names = zin.namelist()
        contents = {}
        for name in names:
            contents[name] = zin.read(name)

    # Modify each slide XML
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'

    for slide_idx, tr_xml in enumerate(transitions_xml):
        slide_name = f'ppt/slides/slide{slide_idx + 1}.xml'
        if slide_name not in contents:
            continue

        slide_xml = contents[slide_name].decode('utf-8')

        # Register namespaces to preserve them
        ET.register_namespace('p', ns_p)
        ET.register_namespace('a', 'http://schemas.openxmlformats.org/drawingml/2006/main')
        ET.register_namespace('r', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships')

        root = ET.fromstring(slide_xml)

        # Remove any existing transition
        existing = root.find(f'{{{ns_p}}}transition')
        if existing is not None:
            root.remove(existing)

        # Parse and append transition element
        tr_elem = ET.fromstring(tr_xml)
        root.append(tr_elem)

        # Re-serialize
        new_xml = ET.tostring(root, encoding='unicode', xml_declaration=False)
        if not new_xml.startswith('<?'):
            new_xml = '<?xml version=\'1.0\' encoding=\'UTF-8\' standalone=\'yes\'?>\n' + new_xml
        contents[slide_name] = new_xml.encode('utf-8')

    # Write modified zip
    with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zout:
        for name, data in contents.items():
            zout.writestr(name, data)

    import os
    os.remove(tmp_path)
    print(f'Transitions added to all 8 slides in {pptx_path}')


create_initial()
