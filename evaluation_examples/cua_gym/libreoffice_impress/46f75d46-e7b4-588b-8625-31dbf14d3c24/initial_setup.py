"""
Initial Setup: Create a 3-slide presentation with left-aligned text on slide 2.
Task ID: impress_tct_087
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_087'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Essay on Renewable Energy"
    slide1.placeholders[1].text = "A Comprehensive Overview of Sustainable Power Sources"

    # --- Slide 2: Content with 3 paragraphs, LEFT-aligned, NO indent ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    txBox = slide2.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8.4), Inches(6.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    paragraphs_text = [
        (
            "Renewable energy has emerged as a critical solution to the growing challenges "
            "of climate change and energy security. Over the past two decades, technological "
            "advances in solar photovoltaic cells, wind turbines, and energy storage systems "
            "have dramatically reduced costs while improving efficiency. Countries around the "
            "world are setting ambitious targets for carbon neutrality, with many aiming to "
            "achieve net-zero emissions by 2050. The transition from fossil fuels to clean "
            "energy sources represents one of the most significant economic and environmental "
            "shifts in modern history."
        ),
        (
            "Solar energy, in particular, has experienced remarkable growth. Global installed "
            "solar capacity surpassed 1,200 gigawatts in 2023, a tenfold increase from just "
            "a decade earlier. The levelized cost of electricity from utility-scale solar "
            "projects has fallen below that of new coal and natural gas plants in most major "
            "markets. Innovations such as perovskite solar cells, bifacial panels, and "
            "floating solar farms continue to push the boundaries of what is achievable. "
            "Meanwhile, distributed rooftop installations empower homeowners and businesses "
            "to generate their own electricity, reducing dependence on centralized grids."
        ),
        (
            "Wind power complements solar as a cornerstone of the renewable energy portfolio. "
            "Offshore wind farms, situated in coastal waters where winds are stronger and more "
            "consistent, are expanding rapidly across Europe, Asia, and North America. Turbine "
            "manufacturers have developed machines with rotor diameters exceeding 200 meters, "
            "capable of generating over 15 megawatts each. Energy storage technologies, "
            "including lithium-ion batteries and emerging alternatives like solid-state and "
            "iron-air batteries, address the intermittency challenge by storing surplus "
            "generation for periods of low wind or sunlight."
        ),
    ]

    # First paragraph uses the default paragraph (index 0)
    p0 = tf.paragraphs[0]
    p0.text = paragraphs_text[0]
    p0.alignment = PP_ALIGN.LEFT
    for run in p0.runs:
        run.font.size = Pt(14)
        run.font.name = "Calibri"
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Second paragraph
    p1 = tf.add_paragraph()
    p1.text = paragraphs_text[1]
    p1.alignment = PP_ALIGN.LEFT
    p1.space_before = Pt(12)
    for run in p1.runs:
        run.font.size = Pt(14)
        run.font.name = "Calibri"
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Third paragraph
    p2 = tf.add_paragraph()
    p2.text = paragraphs_text[2]
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(12)
    for run in p2.runs:
        run.font.size = Pt(14)
        run.font.name = "Calibri"
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 3: Conclusion slide ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    txBox3 = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True

    p3 = tf3.paragraphs[0]
    p3.text = "Key Takeaways"
    p3.alignment = PP_ALIGN.CENTER
    for run in p3.runs:
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.name = "Calibri"
        run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    bullets = [
        "Renewable energy costs have fallen dramatically over the past decade",
        "Solar and wind power now compete directly with fossil fuels on price",
        "Energy storage is the key enabler for reliable clean energy systems",
        "Government policies and private investment drive continued growth",
    ]
    for bullet_text in bullets:
        bp = tf3.add_paragraph()
        bp.text = bullet_text
        bp.alignment = PP_ALIGN.LEFT
        bp.space_before = Pt(8)
        bp.level = 0
        for run in bp.runs:
            run.font.size = Pt(18)
            run.font.name = "Calibri"
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
