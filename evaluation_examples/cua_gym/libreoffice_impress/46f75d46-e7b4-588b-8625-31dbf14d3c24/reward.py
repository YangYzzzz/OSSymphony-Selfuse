"""
Reward Script: Make text on slide 2 justified with 0.5 inch first line indent
Task ID: impress_tct_087
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): All 3 paragraphs have JUSTIFY alignment
  Component 2 (0.5): All 3 paragraphs have first line indent of 0.5 inches (457200 EMU)
"""

import os

from pptx import Presentation
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_087'

# 0.5 inches = 457200 EMU
EXPECTED_INDENT_EMU = 457200
INDENT_TOLERANCE = 5000  # small tolerance for rounding


def persist_app_state(domain):
    """Try to save any unsaved LibreOffice changes via Ctrl+S."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print("PERSIST: ctrl+s sent for %s" % domain)
        except Exception as e:
            print("PERSIST_WARN: save hook failed: %s" % e)


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print("CRITICAL: Cannot load file %s: %s" % (file_path, e))
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 2 slides
    if len(prs.slides) < 2:
        print("PRECONDITION FAIL: Need at least 2 slides, found %d" % len(prs.slides))
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[1]  # Slide 2 (0-indexed)

    # Find the text box with 3 paragraphs on slide 2
    target_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame:
            # Look for text box with substantive content (3 paragraphs with text)
            paras_with_text = [p for p in shape.text_frame.paragraphs if p.text.strip()]
            if len(paras_with_text) >= 3:
                target_shape = shape
                break

    if target_shape is None:
        print("PRECONDITION FAIL: No text box with 3+ paragraphs found on slide 2")
        print("REWARD: 0.0")
        return 0.0

    paragraphs = [p for p in target_shape.text_frame.paragraphs if p.text.strip()]
    print("Found text box with %d non-empty paragraphs on slide 2" % len(paragraphs))

    # Component 1: All 3 paragraphs have JUSTIFY alignment (0.5 points)
    try:
        justified_count = 0
        for i, para in enumerate(paragraphs[:3]):
            align = para.alignment
            if align == PP_ALIGN.JUSTIFY:
                justified_count += 1
                print("PASS: Para %d alignment is JUSTIFY" % i)
            else:
                print("FAIL: Para %d alignment is %s, expected JUSTIFY" % (i, align))

        if justified_count == 3:
            print("PASS: Component 1 -- All 3 paragraphs justified (0.5 pts)")
            total_score += 0.5
        elif justified_count > 0:
            partial = round(0.5 * justified_count / 3, 2)
            print("PARTIAL: Component 1 -- %d/3 paragraphs justified (%.2f pts)" % (justified_count, partial))
            total_score += partial
        else:
            print("FAIL: Component 1 -- No paragraphs are justified")
    except Exception as e:
        print("ERROR: Component 1 -- %s" % e)

    # Component 2: All 3 paragraphs have first line indent of 0.5 inches (0.5 points)
    try:
        ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        indent_count = 0
        for i, para in enumerate(paragraphs[:3]):
            pPr = para._p.find(ns + 'pPr')
            indent_val = None
            if pPr is not None:
                indent_str = pPr.get('indent')
                if indent_str is not None:
                    indent_val = int(indent_str)

            if indent_val is not None and abs(indent_val - EXPECTED_INDENT_EMU) <= INDENT_TOLERANCE:
                indent_count += 1
                print("PASS: Para %d first line indent = %d EMU (expected ~%d)" % (i, indent_val, EXPECTED_INDENT_EMU))
            else:
                print("FAIL: Para %d first line indent = %s EMU, expected ~%d" % (i, indent_val, EXPECTED_INDENT_EMU))

        if indent_count == 3:
            print("PASS: Component 2 -- All 3 paragraphs have 0.5in indent (0.5 pts)")
            total_score += 0.5
        elif indent_count > 0:
            partial = round(0.5 * indent_count / 3, 2)
            print("PARTIAL: Component 2 -- %d/3 paragraphs have correct indent (%.2f pts)" % (indent_count, partial))
            total_score += partial
        else:
            print("FAIL: Component 2 -- No paragraphs have correct first line indent")
    except Exception as e:
        print("ERROR: Component 2 -- %s" % e)

    final_score = min(total_score, 1.0)
    print("\nScore: %.2f/1.0" % total_score)
    print("REWARD: %.1f" % final_score)
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = '%s/%s.pptx' % (WORKDIR, TASK_ID)
if not os.path.exists(file_path):
    print("File not found: %s" % file_path)
    print("REWARD: 0.0")
else:
    verify_task(file_path)
