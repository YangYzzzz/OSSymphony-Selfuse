"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 46 I only want to touch the first three bullet points—can I turn those into a numbered list that starts at 3 instead of the usual 1, 2, 3? Using LibreOffice Impress and hoping there’s a quick way to do just that.
Generated: 2025-09-10 22:18:48
Status: success
Model: azure-o3
Total Steps: 4
"""

from pptx import Presentation
import os
import re


def verify_numbered_list_starting_at_3(file_path: str, slide_number: int = 46) -> float:
    """Reward script for the LibreOffice Impress task.

    Task requirements (Slide 46):
    1. Only the first three bullet points must be converted to a numbered list.
    2. That numbered list must start at 3 (so the first, second, third items are numbered 3, 4, 5).
    3. All remaining bullets on the slide must stay un-numbered.

    Progressive scoring (adds up to 1.0):
        0.0  – critical failure / file missing / slide missing
        +0.20 – first three bullet texts are non-empty
        +0.30 – first three bullets are **actually numbered** (buAutoNum present)
        +0.30 – numbering sequence is exactly 3, 4, 5
        +0.20 – all subsequent bullets are **not** numbered
        ——
        1.00 – perfect fulfilment of all requirements
    """
    max_score = 1.0
    score = 0.0

    # 1. Load the presentation (no points just for loading)
    if not os.path.exists(file_path):
        print(f"✗ File does not exist: {file_path}")
        return 0.0
    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded presentation with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Failed to load PPTX: {e}")
        return 0.0

    # 2. Ensure the requested slide exists
    target_index = slide_number - 1  # zero-based index
    if len(prs.slides) <= target_index:
        print(f"✗ Presentation contains only {len(prs.slides)} slides, expected ≥ {slide_number}")
        return 0.0
    slide = prs.slides[target_index]
    print(f"✓ Slide {slide_number} found")

    # 3. Find a shape with multiple paragraphs (bullet list)
    list_shape = None
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text and "\n" in shape.text:
            # Likely candidate: a text frame containing several bullet paragraphs
            list_shape = shape
            break
    if list_shape is None:
        print("✗ Could not locate a multi-paragraph text shape on the target slide.")
        return 0.0

    text_frame = list_shape.text_frame
    paragraphs = text_frame.paragraphs
    para_count = len(paragraphs)
    print(f"✓ Found text frame with {para_count} paragraphs")

    # The slide should contain at least 5 bullets (3 numbered + 2 unnumbered)
    if para_count < 5:
        print("✗ Expected at least 5 bullet paragraphs, found", para_count)
        return 0.0

    # 4. Verify the first three bullet texts are present (non-empty)
    first_three_non_empty = all(paragraphs[i].text.strip() for i in range(3))
    if first_three_non_empty:
        score += 0.20
        print("✓ First three bullet texts are non-empty (+0.20)")
    else:
        print("✗ One of the first three bullet texts is empty")

    # 5. Check each of the first three paragraphs is numbered and capture startAt values
    numbered_flags = []
    start_numbers = []
    for i in range(3):
        p_xml = paragraphs[i]._p.xml
        match = re.search(r"<a:buAutoNum[^>]*?startAt=\"(\d+)\"", p_xml)
        if match:
            numbered_flags.append(True)
            start_numbers.append(int(match.group(1)))
        else:
            numbered_flags.append(False)
            start_numbers.append(None)
    if all(numbered_flags):
        score += 0.30
        print("✓ First three bullets are numbered (+0.30)")
    else:
        print("✗ The first three bullets are not all numbered")

    # 6. Validate the numbering sequence 3, 4, 5.
    if start_numbers == [3, 4, 5]:
        score += 0.30
        print("✓ Numbering sequence is exactly 3,4,5 (+0.30)")
    else:
        print(f"✗ Incorrect numbering sequence: {start_numbers}")

    # 7. Ensure all remaining paragraphs are NOT numbered
    remainder_correct = True
    for idx in range(3, para_count):
        if "<a:buAutoNum" in paragraphs[idx]._p.xml:
            remainder_correct = False
            print(f"✗ Paragraph {idx} (zero-based) is numbered but should not be")
            break
    if remainder_correct:
        score += 0.20
        print("✓ Remaining paragraphs are unnumbered (+0.20)")

    final_score = min(score, max_score)
    print(f"REWARD: {final_score}")
    return final_score


# Execute verification when run as a script
if __name__ == "__main__":
    FILE_PATH = "/home/user/on_slide_46_i_only_want_to_touch_the_first_three_bullet_pointscan_i_turn_those_into_a_numbered_list__golden.pptx"
    verify_numbered_list_starting_at_3(FILE_PATH)

