"""
Reward Script: Add presentation title text box to slide master in bottom-left corner
Task ID: impress_gf1_040
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): New text box on slide master contains presentation title text
  Component 2 (0.4): Text formatting — 10pt italic grey #757575
  Component 3 (0.3): Text box positioned in bottom-left corner of slide
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_gf1_040'


def persist_app_state(domain):
    """Save any unsaved LibreOffice changes before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Pt, Emu
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Get the presentation title from core properties
    pres_title = prs.core_properties.title or ""
    print(f"INFO: Presentation title from core properties: {repr(pres_title)}")

    # Get slide dimensions for position checks
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    print(f"INFO: Slide dimensions: {slide_w} x {slide_h}")

    # Find the slide master
    if len(prs.slide_masters) == 0:
        print("CRITICAL: No slide masters found")
        print("REWARD: 0.0")
        return 0.0

    master = prs.slide_masters[0]

    # Find text boxes on the slide master that contain the presentation title
    # We look for TEXT_BOX shapes (type 17) that are NOT standard placeholders
    # and contain text matching the presentation title
    title_textboxes = []
    for shape in master.shapes:
        if not shape.has_text_frame:
            continue
        # Check if this is a text box (not a placeholder) or any shape with matching text
        shape_text = shape.text_frame.text.strip()
        # Match: the text should contain the presentation title
        if pres_title and pres_title.strip() and pres_title.strip().lower() in shape_text.lower():
            # Exclude standard placeholder shapes (Title Placeholder, Text Placeholder, etc.)
            # by checking shape type or name patterns
            is_standard_placeholder = any(kw in shape.name for kw in [
                'Title Placeholder', 'Text Placeholder', 'Date Placeholder',
                'Footer Placeholder', 'Slide Number Placeholder'
            ])
            if not is_standard_placeholder:
                title_textboxes.append(shape)
                print(f"INFO: Found candidate shape: name={shape.name}, type={shape.shape_type}, text={repr(shape_text)}")

    # Component 1: Text box on slide master contains presentation title (0.3 points)
    try:
        if len(title_textboxes) > 0:
            print(f"PASS: Component 1 — Found {len(title_textboxes)} text box(es) with presentation title on slide master (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 — No text box with presentation title '{pres_title}' found on slide master")
            # If no title textbox found, remaining checks can't pass
            final_score = min(total_score, 1.0)
            print(f"\nScore: {total_score}/1.0")
            print(f"REWARD: {final_score}")
            return final_score
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Use the first matching text box for further checks
    tb = title_textboxes[0]

    # Component 2: Text formatting — 10pt, italic, grey #757575 (0.4 points)
    # Sub-checks: size (0.15), italic (0.1), color (0.15)
    try:
        format_score = 0.0
        # Get runs from the text frame
        all_runs = []
        for para in tb.text_frame.paragraphs:
            for run in para.runs:
                if (run.text or "").strip():
                    all_runs.append(run)

        if len(all_runs) == 0:
            print("FAIL: Component 2 — No text runs found in the title text box")
        else:
            # Check first non-empty run (primary formatting)
            run = all_runs[0]

            # Sub-check 2a: Font size == 10pt (127000 EMU)
            font_size = run.font.size
            if font_size is not None and abs(font_size - 127000) < 5000:
                print(f"PASS: Component 2a — Font size is {font_size} EMU (~{font_size/12700:.1f}pt), expected 10pt (0.15 pts)")
                format_score += 0.15
            else:
                print(f"FAIL: Component 2a — Font size is {font_size}, expected 127000 EMU (10pt)")

            # Sub-check 2b: Italic is True
            is_italic = run.font.italic
            if is_italic is True:
                print(f"PASS: Component 2b — Text is italic (0.1 pts)")
                format_score += 0.1
            else:
                print(f"FAIL: Component 2b — Italic is {is_italic}, expected True")

            # Sub-check 2c: Color is #757575
            try:
                color_rgb = None
                if run.font.color.type is not None:
                    color_rgb = str(run.font.color.rgb)
                if color_rgb and color_rgb.upper() == "757575":
                    print(f"PASS: Component 2c — Color is #{color_rgb} (0.15 pts)")
                    format_score += 0.15
                else:
                    print(f"FAIL: Component 2c — Color is {color_rgb}, expected 757575")
            except Exception as ce:
                print(f"FAIL: Component 2c — Could not read color: {ce}")

        total_score += format_score
        if format_score >= 0.35:
            print(f"PASS: Component 2 — Formatting mostly correct ({format_score:.2f}/0.4 pts)")
        else:
            print(f"PARTIAL: Component 2 — Formatting score: {format_score:.2f}/0.4")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Position — text box is in bottom-left corner (0.3 points)
    try:
        shape_left = tb.left
        shape_top = tb.top
        shape_right = shape_left + tb.width
        shape_bottom = shape_top + tb.height

        # Bottom-left criteria:
        # - Left edge should be in the left quarter of the slide
        # - Top/bottom should be in the bottom quarter of the slide
        is_left = shape_left < slide_w * 0.35
        is_bottom = shape_top > slide_h * 0.75

        print(f"INFO: Shape position — left={shape_left}, top={shape_top}, right={shape_right}, bottom={shape_bottom}")
        print(f"INFO: Left check: {shape_left} < {slide_w * 0.35:.0f} = {is_left}")
        print(f"INFO: Bottom check: {shape_top} > {slide_h * 0.75:.0f} = {is_bottom}")

        if is_left and is_bottom:
            print(f"PASS: Component 3 — Text box is in bottom-left corner (0.3 pts)")
            total_score += 0.3
        elif is_left or is_bottom:
            print(f"PARTIAL: Component 3 — Partially in bottom-left (left={is_left}, bottom={is_bottom}) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 3 — Text box is not in bottom-left corner")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
