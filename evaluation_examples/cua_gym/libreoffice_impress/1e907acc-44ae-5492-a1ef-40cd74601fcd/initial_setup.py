"""
Initial Setup: Create a 15-slide presentation with presentation title set in properties.
Task ID: impress_gf1_040
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf1_040'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # Set presentation title in core properties
    prs.core_properties.title = 'Annual Strategy Review 2024'

    # Slide dimensions (standard widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide content data for 15 slides
    slide_content = [
        {
            'layout': 0,  # Title Slide
            'title': 'Annual Strategy Review 2024',
            'subtitle': 'Prepared by the Executive Strategy Team\nQ4 Board Presentation'
        },
        {
            'layout': 1,  # Title + Content
            'title': 'Executive Summary',
            'body': 'Revenue grew 18% year-over-year to $4.2B\nMarket share expanded in 3 key regions\nNew product launches exceeded targets by 12%\nCustomer retention rate improved to 94.7%'
        },
        {
            'layout': 1,
            'title': 'Financial Highlights',
            'body': 'Total Revenue: $4.2B (+18% YoY)\nGross Margin: 62.3% (up from 59.1%)\nOperating Income: $890M (+22% YoY)\nFree Cash Flow: $1.1B\nR&D Investment: $620M (14.8% of revenue)'
        },
        {
            'layout': 1,
            'title': 'Market Position Analysis',
            'body': 'North America: #1 position maintained (34% share)\nEurope: Moved from #3 to #2 (21% share)\nAsia-Pacific: Fastest growing region (+31% YoY)\nLatin America: New market entry completed'
        },
        {
            'layout': 1,
            'title': 'Product Portfolio Performance',
            'body': 'Enterprise Suite: $1.8B revenue, 156K active licenses\nCloud Platform: $1.2B revenue, 89% renewal rate\nSecurity Solutions: $680M revenue, fastest growing segment\nProfessional Services: $520M revenue, NPS score 72'
        },
        {
            'layout': 1,
            'title': 'Customer Acquisition Metrics',
            'body': 'New enterprise customers: 2,340 (+28%)\nAverage deal size: $185K (up from $142K)\nSales cycle reduced from 90 to 72 days\nPartner-sourced revenue: 38% of total'
        },
        {
            'layout': 1,
            'title': 'Technology Innovation Roadmap',
            'body': 'AI-powered analytics engine launched in Q2\n47 patents filed, 31 granted\nPlatform migration to microservices: 78% complete\nNew developer ecosystem: 12,000+ registered developers'
        },
        {
            'layout': 1,
            'title': 'Operational Excellence',
            'body': 'System uptime: 99.97% (exceeded SLA)\nCustomer support CSAT: 4.6/5.0\nEmployee engagement score: 82nd percentile\nData center energy efficiency improved 15%'
        },
        {
            'layout': 1,
            'title': 'Talent & Organization',
            'body': 'Headcount: 8,450 employees across 22 offices\nVoluntary turnover: 8.2% (industry avg 14.5%)\nDiversity hiring increased 23% YoY\nInternal promotion rate: 34%'
        },
        {
            'layout': 1,
            'title': 'Strategic Partnerships',
            'body': 'AWS partnership expanded to Platinum tier\nMicrosoft co-sell agreement renewed\nNew alliance with SAP for ERP integration\nConsulting partnerships: Deloitte, Accenture, PwC'
        },
        {
            'layout': 1,
            'title': 'Risk Assessment',
            'body': 'Regulatory compliance: SOC 2 Type II achieved\nCybersecurity: Zero major incidents in 2024\nSupply chain: Diversified to 3 cloud providers\nGeopolitical: Monitoring APAC regulatory changes'
        },
        {
            'layout': 1,
            'title': 'Sustainability Initiatives',
            'body': 'Carbon neutral operations by 2026 (on track)\nRenewable energy usage: 72% of total\nPaper-free offices: 18 of 22 locations\nCommunity investment: $8.5M in STEM education'
        },
        {
            'layout': 1,
            'title': 'Competitive Landscape',
            'body': 'Primary competitor revenue growth: 11% (vs our 18%)\nWin rate in head-to-head deals: 62%\nAnalyst rankings: Leader in 4 Gartner quadrants\nBrand awareness increased 15 points'
        },
        {
            'layout': 1,
            'title': '2025 Strategic Priorities',
            'body': 'Accelerate AI/ML product integration\nExpand APAC presence with 3 new offices\nLaunch vertical-specific solutions for Healthcare and Finance\nAchieve $5B revenue milestone\nComplete platform modernization'
        },
        {
            'layout': 1,
            'title': 'Thank You & Questions',
            'body': 'Contact: strategy@company.com\nNext review: Q1 2025 Board Meeting\nDetailed appendix available in the shared drive\nPlease submit questions via the board portal'
        },
    ]

    for i, content in enumerate(slide_content):
        layout_idx = content['layout']
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = content['title']

        # Set body/subtitle
        body_key = 'subtitle' if layout_idx == 0 else 'body'
        if body_key in content:
            # Find the body placeholder
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:  # subtitle or body
                    tf = ph.text_frame
                    lines = content[body_key].split('\n')
                    tf.paragraphs[0].text = lines[0]
                    for line in lines[1:]:
                        p = tf.add_paragraph()
                        p.text = line

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Slides: {len(prs.slides)}')
    print(f'Presentation title: {prs.core_properties.title}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
