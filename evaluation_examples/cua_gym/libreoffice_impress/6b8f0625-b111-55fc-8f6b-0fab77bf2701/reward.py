"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 33 of my LibreOffice Impress deck, I need to add a horizontal line that is exactly 8 cm long and have it sit perfectly flush—same left margin—right under the title text box. I keep missing the alignment by a hair. What’s the quickest way to snap that line into the correct spot?
Generated: 2025-09-10 13:14:45
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
from pptx import Presentation

# Constants
EMUS_PER_CM = 360000  # EMU (English Metric Unit) to centimetres conversion factor


def cm(emus: int) -> float:
    """Convert EMUs to centimetres."""
    return emus / EMUS_PER_CM


def verify_line_alignment(pptx_path: str) -> float:
    """Verify that slide 33 contains an 8 cm horizontal line perfectly
    aligned (same left margin) directly under the title text box.

    Scoring (progressive):
        • 0.4 pts – Horizontal line of 8 cm ± 0.1 cm exists
        • 0.3 pts – Line left margin aligns with title left margin (± 0.05 cm)
        • 0.3 pts – Line placed flush under title (gap 0-0.5 cm)
    Returns a float between 0.0 and 1.0.
    """

    print(f"Verifying line-alignment task for file: {pptx_path}\n")

    # ---------- Basic file checks (no points awarded) ----------
    if not os.path.exists(pptx_path):
        print("✗ File does not exist")
        return 0.0
    try:
        prs = Presentation(pptx_path)
    except Exception as exc:
        print(f"✗ Unable to open presentation: {exc}")
        return 0.0

    expected_slide_idx = 32  # zero-based index for slide 33
    if len(prs.slides) <= expected_slide_idx:
        print(f"✗ Presentation has only {len(prs.slides)} slide(s); slide 33 missing")
        return 0.0

    slide = prs.slides[expected_slide_idx]
    print("✓ Slide 33 located")

    # ---------- Locate title placeholder ----------
    title_shape = None
    for shp in slide.shapes:
        if getattr(shp, "name", "").lower().startswith("title"):
            title_shape = shp
            break
    if not title_shape:
        print("✗ Title placeholder not found on slide 33")
        return 0.0

    title_left = title_shape.left
    title_bottom = title_shape.top + title_shape.height
    print(
        f"Title left = {cm(title_left):.2f} cm, bottom = {cm(title_bottom):.2f} cm"
    )

    # ---------- Find horizontal line candidates ----------
    # shape_type 6 = LINE, 9 = CONNECTOR (often used for straight lines)
    line_candidates = [
        shp for shp in slide.shapes if shp.shape_type in (6, 9)
    ]
    print(f"Found {len(line_candidates)} potential line shape(s)")

    ideal_width = 8 * EMUS_PER_CM
    width_tol = 0.1 * EMUS_PER_CM   # ±0.1 cm
    pos_tol = 0.05 * EMUS_PER_CM    # ±0.05 cm for alignment
    vertical_max_gap_cm = 0.5        # ≤0.5 cm below title

    chosen_line = None
    for line in line_candidates:
        correct_length = abs(line.width - ideal_width) <= width_tol
        horizontal = line.height <= pos_tol  # virtually zero height
        if correct_length and horizontal:
            chosen_line = line
            break

    score = 0.0

    # ---------- Scoring: length ----------
    if not chosen_line:
        print("✗ No 8 cm horizontal line found (±0.1 cm tolerance)")
        return score  # 0.0

    print("✓ 8 cm horizontal line located (0.4 pts)")
    score += 0.4

    # ---------- Scoring: left alignment ----------
    if abs(chosen_line.left - title_left) <= pos_tol:
        print("✓ Line left margin aligns with title (0.3 pts)")
        score += 0.3
    else:
        delta_left = cm(abs(chosen_line.left - title_left))
        print(f"✗ Left alignment off by {delta_left:.2f} cm (>0.05 cm)")

    # ---------- Scoring: vertical placement ----------
    gap_cm = cm(chosen_line.top) - cm(title_bottom)
    if 0.0 <= gap_cm <= vertical_max_gap_cm:
        print(f"✓ Line positioned {gap_cm:.2f} cm below title (0.3 pts)")
        score += 0.3
    else:
        print(
            f"✗ Vertical gap is {gap_cm:.2f} cm (allowed 0–{vertical_max_gap_cm} cm)"
        )

    final_score = round(min(score, 1.0), 2)
    print(f"\nTotal score: {final_score}")
    return final_score


if __name__ == "__main__":
    FILE_PATH = (
        "/home/user/"
        "on_slide_33_of_my_libreoffice_impress_deck_i_need_to_add_a_horizontal_"
        "line_that_is_exactly_8_cm_long_golden.pptx"
    )
    reward = verify_line_alignment(FILE_PATH)
    print(f"REWARD: {reward}")
