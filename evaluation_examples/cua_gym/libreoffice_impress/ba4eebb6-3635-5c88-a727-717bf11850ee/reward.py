"""
Reward Script: Move the image on slide 3 to the center of the slide.
Task ID: osworld_impress_image_fill_slide_009
Domain: libreoffice_impress
Scoring:
    Component 1 (0.5): Image on slide 3 is horizontally centered (center_x ≈ slide_width/2)
    Component 2 (0.5): Image on slide 3 is vertically centered (center_y ≈ slide_height/2)
    Total: 1.0
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_image_fill_slide_009'

# Tolerance for centering check: 0.5% relative tolerance (from SKILL.md)
TOLERANCE = 0.005


def is_approximately_equal(val1, val2, tolerance=TOLERANCE):
    """Check if two values are approximately equal with relative tolerance."""
    if val1 == val2:
        return True
    if val1 == 0 or val2 == 0:
        return val1 == val2
    return abs(val1 - val2) / max(abs(val1), abs(val2)) <= tolerance


def verify_task(file_path):
    """
    Verify that the image on slide 3 has been moved to the center of the slide.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"CRITICAL: Presentation has only {len(prs.slides)} slides; expected at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide_width = prs.slide_width    # 9144000 EMU
    slide_height = prs.slide_height  # 6858000 EMU
    slide_center_x = slide_width / 2   # 4572000 EMU
    slide_center_y = slide_height / 2  # 3429000 EMU

    print(f"Slide dimensions: width={slide_width}, height={slide_height}")
    print(f"Slide center: ({slide_center_x}, {slide_center_y})")

    slide3 = prs.slides[2]  # 0-indexed: slide index 2 = slide 3

    # Find the picture shape on slide 3
    picture_shape = None
    for shape in slide3.shapes:
        if shape.shape_type == 13:  # PICTURE type = 13
            picture_shape = shape
            break

    if picture_shape is None:
        print("FAIL: No picture/image found on slide 3")
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    shape_center_x = picture_shape.left + picture_shape.width / 2
    shape_center_y = picture_shape.top + picture_shape.height / 2
    print(f"Picture shape: left={picture_shape.left}, top={picture_shape.top}, "
          f"width={picture_shape.width}, height={picture_shape.height}")
    print(f"Shape center: ({shape_center_x}, {shape_center_y})")

    # Component 1: Image is horizontally centered (center_x ≈ slide_width/2) (0.5 points)
    # This FAILS on initial (center_x=7040880, far right) and PASSES on golden (center_x=4572000)
    try:
        if is_approximately_equal(shape_center_x, slide_center_x, tolerance=0.02):
            print(f"PASS: Component 1 — Image horizontally centered "
                  f"(shape_center_x={shape_center_x}, slide_center_x={slide_center_x}) (0.5 pts)")
            total_score += 0.5
        else:
            offset_x = shape_center_x - slide_center_x
            print(f"FAIL: Component 1 — Image NOT horizontally centered; "
                  f"shape_center_x={shape_center_x}, slide_center_x={slide_center_x}, "
                  f"offset={offset_x}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Image is vertically centered (center_y ≈ slide_height/2) (0.5 points)
    # This FAILS on initial (center_y=5212080, lower area) and PASSES on golden (center_y=3429000)
    try:
        if is_approximately_equal(shape_center_y, slide_center_y, tolerance=0.02):
            print(f"PASS: Component 2 — Image vertically centered "
                  f"(shape_center_y={shape_center_y}, slide_center_y={slide_center_y}) (0.5 pts)")
            total_score += 0.5
        else:
            offset_y = shape_center_y - slide_center_y
            print(f"FAIL: Component 2 — Image NOT vertically centered; "
                  f"shape_center_y={shape_center_y}, slide_center_y={slide_center_y}, "
                  f"offset={offset_y}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the VM environment
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
