"""
Initial Setup: 5-slide product catalog with product image in lower-right of slide 3.
Task ID: osworld_impress_image_fill_slide_009
Domain: libreoffice_impress
"""

import os
import io
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_image_fill_slide_009'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
IMG_PATH = f'{WORKDIR}/{TASK_ID}_product.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_product_image(path: str):
    """Create a realistic product photo (a stylized laptop/tech product image)."""
    width, height = 400, 300
    img = Image.new('RGB', (width, height), color=(245, 245, 250))
    draw = ImageDraw.Draw(img)

    # Draw laptop body
    draw.rectangle([60, 80, 340, 220], fill=(50, 50, 60), outline=(30, 30, 40), width=3)
    # Screen area
    draw.rectangle([75, 90, 325, 205], fill=(20, 120, 200))
    # Screen reflection glare
    draw.polygon([(80, 92), (160, 92), (120, 140)], fill=(100, 180, 255, 80))
    # Screen text simulation
    draw.rectangle([90, 110, 310, 125], fill=(255, 255, 255, 120))
    draw.rectangle([90, 130, 260, 140], fill=(200, 230, 255, 100))
    draw.rectangle([90, 148, 290, 155], fill=(200, 230, 255, 100))
    draw.rectangle([90, 163, 240, 170], fill=(200, 230, 255, 100))
    # Keyboard area
    draw.rectangle([70, 222, 330, 245], fill=(60, 60, 70), outline=(40, 40, 50), width=2)
    # Keyboard keys (simplified)
    for row in range(2):
        for col in range(12):
            kx = 80 + col * 21
            ky = 226 + row * 9
            draw.rectangle([kx, ky, kx+18, ky+6], fill=(80, 80, 90), outline=(50, 50, 60))
    # Trackpad
    draw.rectangle([155, 248, 245, 260], fill=(70, 70, 80), outline=(50, 50, 60))

    # Product label
    draw.rectangle([0, 260, width, height], fill=(230, 240, 255))
    draw.text((width // 2 - 80, 265), "TechPro X15 Laptop", fill=(30, 30, 100))
    draw.text((width // 2 - 50, 280), "Model: TPX15-2025", fill=(80, 80, 120))

    img.save(path, 'PNG')
    print(f'Product image created: {path}')


def create_initial():
    create_product_image(IMG_PATH)

    prs = Presentation()
    # Standard widescreen: 10 x 7.5 inches
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "TechPro Innovations"
    slide1.placeholders[1].text = "2025 Product Catalog\nBringing Technology to Life"
    # Style title
    for run in slide1.shapes.title.text_frame.paragraphs[0].runs:
        run.font.bold = True
        run.font.size = Pt(40)
        run.font.color.rgb = RGBColor(0x1A, 0x3A, 0x6C)
    # Background
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFF)

    # ---- Slide 2: Product Overview ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Our Product Lineup"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "TechPro X15 Laptop"
    items2 = [
        "UltraSound Pro Headphones",
        "SmartView Monitor 4K",
        "PowerHub 10-Port USB Dock",
        "QuickCharge Wireless Pad",
    ]
    for item in items2:
        p = tf2.add_paragraph()
        p.text = item
        p.level = 1
    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFF)

    # ---- Slide 3: Featured Product (image in lower-right, NOT centered) ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Title text box
    title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1.0))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "TechPro X15 Laptop"
    p_title.alignment = PP_ALIGN.LEFT
    run_title = p_title.runs[0]
    run_title.font.bold = True
    run_title.font.size = Pt(32)
    run_title.font.color.rgb = RGBColor(0x1A, 0x3A, 0x6C)

    # Description text box
    desc_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.0), Inches(4.5))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = "Key Features:"
    run_d = p_desc.runs[0]
    run_d.font.bold = True
    run_d.font.size = Pt(18)
    run_d.font.color.rgb = RGBColor(0x2C, 0x2C, 0x2C)

    specs = [
        "Intel Core i9-14900H Processor",
        "32GB DDR5 RAM",
        "1TB NVMe SSD",
        "15.6\" 4K OLED Display",
        "NVIDIA RTX 4070 Graphics",
        "20-hour battery life",
        "Price: $1,899.99",
    ]
    for spec in specs:
        p = tf_desc.add_paragraph()
        p.text = f"• {spec}"
        p.level = 0
        if p.runs:
            p.runs[0].font.size = Pt(14)
            p.runs[0].font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # Product image in LOWER-RIGHT area (NOT centered)
    img_w = Inches(4.0)
    img_h = Inches(3.0)
    # Lower-right: leave ~0.3 inch margin from right and bottom
    img_left = slide_width - img_w - Inches(0.3)
    img_top = slide_height - img_h - Inches(0.3)
    slide3.shapes.add_picture(IMG_PATH, img_left, img_top, img_w, img_h)

    # Background
    fill3 = slide3.background.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # ---- Slide 4: Accessories ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Top Accessories"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "UltraSound Pro Headphones — $299.99"
    accessories = [
        "SmartView Monitor 4K — $649.99",
        "PowerHub 10-Port USB Dock — $129.99",
        "QuickCharge Wireless Pad — $59.99",
        "TechPro Laptop Sleeve — $39.99",
        "ErgoMouse Precision — $89.99",
    ]
    for acc in accessories:
        p = tf4.add_paragraph()
        p.text = acc
        p.level = 1
    fill4 = slide4.background.fill
    fill4.solid()
    fill4.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFF)

    # ---- Slide 5: Contact / CTA ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[0])
    slide5.shapes.title.text = "Order Today"
    slide5.placeholders[1].text = (
        "Visit: www.techpro-innovations.com\n"
        "Call: 1-800-TECHPRO\n"
        "Email: sales@techpro-innovations.com"
    )
    fill5 = slide5.background.fill
    fill5.solid()
    fill5.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x6C)
    for run in slide5.shapes.title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.bold = True
        run.font.size = Pt(36)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
