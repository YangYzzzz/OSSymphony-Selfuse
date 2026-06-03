"""
Reward Script: Change the font color of all text on slide 3 to dark red (#8B0000)
Task ID: osworld_impress_underline_darkred_table_002
Domain: libreoffice_impress
Scoring:
  Component 1 — Title text on slide 3 is dark red (#8B0000): 0.4 points
  Component 2 — All body/content text on slide 3 is dark red (#8B0000): 0.6 points
  Total: 1.0
"""

import os
from pptx import Presentation
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_underline_darkred_table_002'

TARGET_COLOR = '8B0000'  # dark red


def get_all_runs_in_shape(shape):
    """Recursively collect all non-empty runs from a shape (handles groups)."""
    runs = []
    if hasattr(shape, 'text_frame') and shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if (run.text or '').strip():
                    runs.append(run)
    if hasattr(shape, 'shapes'):
        for sub in shape.shapes:
            runs.extend(get_all_runs_in_shape(sub))
    return runs


def run_has_dark_red_color(run):
    """Check if a run has the exact dark red color #8B0000."""
    try:
        if run.font.color.type is None:
            return False
        rgb = run.font.color.rgb
        return str(rgb).upper() == TARGET_COLOR.upper()
    except Exception:
        return False


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"CRITICAL: Expected at least 3 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[2]  # slide 3 is index 2

    # Component 1: Title text on slide 3 is dark red (#8B0000) (0.4 points)
    # The title shape is the first placeholder (shape index 0 with name "Title 1")
    # Task: all text on slide 3 should be changed from black (000000) to dark red (8B0000)
    try:
        title_shape = None
        for shape in slide.shapes:
            if shape.has_text_frame and 'Title' in shape.name:
                title_shape = shape
                break

        if title_shape is None:
            print("FAIL: Component 1 — Title shape not found on slide 3")
        else:
            title_runs = get_all_runs_in_shape(title_shape)
            if not title_runs:
                print("FAIL: Component 1 — No runs found in title shape")
            else:
                all_dark_red = all(run_has_dark_red_color(r) for r in title_runs)
                if all_dark_red:
                    print(f"PASS: Component 1 — Title text on slide 3 is dark red (#8B0000) ({len(title_runs)} run(s) checked) (0.4 pts)")
                    total_score += 0.4
                else:
                    failed_runs = [(r.text[:20], str(r.font.color.rgb) if r.font.color.type is not None else 'inherited')
                                   for r in title_runs if not run_has_dark_red_color(r)]
                    print(f"FAIL: Component 1 — Title text not all dark red. Failures: {failed_runs}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: All body/content text on slide 3 is dark red (#8B0000) (0.6 points)
    # Body text is the second placeholder (Content Placeholder 2)
    # Each of the 5 bullet lines must have dark red color
    try:
        body_shape = None
        for shape in slide.shapes:
            if shape.has_text_frame and 'Title' not in shape.name:
                body_shape = shape
                break

        if body_shape is None:
            print("FAIL: Component 2 — Body/content shape not found on slide 3")
        else:
            body_runs = get_all_runs_in_shape(body_shape)
            if not body_runs:
                print("FAIL: Component 2 — No runs found in body shape")
            else:
                all_dark_red = all(run_has_dark_red_color(r) for r in body_runs)
                if all_dark_red:
                    print(f"PASS: Component 2 — All body text on slide 3 is dark red (#8B0000) ({len(body_runs)} run(s) checked) (0.6 pts)")
                    total_score += 0.6
                else:
                    failed_runs = [(r.text[:20], str(r.font.color.rgb) if r.font.color.type is not None else 'inherited')
                                   for r in body_runs if not run_has_dark_red_color(r)]
                    print(f"FAIL: Component 2 — Body text not all dark red. Failures (sample): {failed_runs[:3]}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
