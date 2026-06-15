"""
Initial Setup: Risk Assessment Presentation - 6 slides with black text on slide 3
Task ID: osworld_impress_underline_darkred_table_002
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_underline_darkred_table_002'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_run_color_black(run):
    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)


def create_initial():
    prs = Presentation()

    # Set slide size (standard widescreen)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    BLACK = RGBColor(0x00, 0x00, 0x00)
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Enterprise Risk Assessment Report"
    sub = slide1.placeholders[1]
    sub.text = "Q2 2025 | Prepared by the Risk Management Division"
    for para in slide1.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
            run.font.size = Pt(36)
    for para in sub.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = DARK_GRAY
            run.font.size = Pt(18)

    # ---- Slide 2: Executive Summary ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    for para in slide2.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
            run.font.size = Pt(32)
    body2 = slide2.placeholders[1]
    body2.text = (
        "This report summarizes identified operational, financial, and strategic risks "
        "for the fiscal quarter ending June 30, 2025.\n"
        "Key findings indicate elevated exposure in three primary risk categories.\n"
        "Mitigation strategies have been outlined for each critical risk area.\n"
        "Board-level review is recommended for risks rated High or Critical."
    )
    for para in body2.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = DARK_GRAY
            run.font.size = Pt(14)

    # ---- Slide 3: Risk Categories (title + body in BLACK — task target) ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Risk Categories Overview"
    for para in slide3.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = BLACK
            run.font.size = Pt(32)
            run.font.bold = True
    body3 = slide3.placeholders[1]
    body3.text = (
        "Operational Risk: Process failures, system outages, and supply chain disruptions.\n"
        "Financial Risk: Currency fluctuation, credit exposure, and liquidity constraints.\n"
        "Strategic Risk: Market shifts, competitive threats, and regulatory changes.\n"
        "Compliance Risk: Data privacy regulations, environmental standards, and labor laws.\n"
        "Reputational Risk: Social media incidents, product recalls, and stakeholder relations."
    )
    for para in body3.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = BLACK
            run.font.size = Pt(14)

    # ---- Slide 4: Operational Risk Details ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Operational Risk Details"
    for para in slide4.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
            run.font.size = Pt(32)
    body4 = slide4.placeholders[1]
    body4.text = (
        "Risk ID: OPS-2025-01 — IT Infrastructure Downtime (Severity: High)\n"
        "Risk ID: OPS-2025-02 — Key Supplier Bankruptcy (Severity: Critical)\n"
        "Risk ID: OPS-2025-03 — Data Breach / Cybersecurity Incident (Severity: Critical)\n"
        "Risk ID: OPS-2025-04 — Employee Turnover in Core Teams (Severity: Medium)\n"
        "Mitigation: Redundant systems, supplier diversification, security audits."
    )
    for para in body4.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = DARK_GRAY
            run.font.size = Pt(13)

    # ---- Slide 5: Financial Risk Details ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Financial Risk Details"
    for para in slide5.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
            run.font.size = Pt(32)
    body5 = slide5.placeholders[1]
    body5.text = (
        "Risk ID: FIN-2025-01 — USD/EUR Exchange Rate Volatility (Severity: Medium)\n"
        "Risk ID: FIN-2025-02 — Client Default on Major Contracts (Severity: High)\n"
        "Risk ID: FIN-2025-03 — Rising Interest Rates Impact on Debt (Severity: Medium)\n"
        "Risk ID: FIN-2025-04 — Reduced Q3 Revenue Forecast (Severity: High)\n"
        "Mitigation: Hedging strategies, diversified client portfolio, cost controls."
    )
    for para in body5.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = DARK_GRAY
            run.font.size = Pt(13)

    # ---- Slide 6: Recommendations & Next Steps ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Recommendations & Next Steps"
    for para in slide6.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
            run.font.size = Pt(32)
    body6 = slide6.placeholders[1]
    body6.text = (
        "1. Convene Risk Committee meeting by July 15, 2025.\n"
        "2. Assign risk owners for all Critical and High-severity items.\n"
        "3. Update Business Continuity Plan with new supplier contingencies.\n"
        "4. Commission external cybersecurity audit — deadline: August 1, 2025.\n"
        "5. Present updated risk register to the Board of Directors in Q3 review."
    )
    for para in body6.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = DARK_GRAY
            run.font.size = Pt(14)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
