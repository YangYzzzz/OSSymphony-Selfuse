"""
Initial Setup: Product Strategy presentation with 10 slides.
Slide 8 has title 'Go-to-Market Process' but no content shapes.
Task ID: impress_exec_054
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_054'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.text = body_lines[0]
    for line in body_lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
    return slide


def add_title_only_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(prs, "Product Strategy 2025",
                    "Accelerating Growth Through Innovation\nPrepared by Strategic Planning Division")

    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "Revenue target: $42M ARR by Q4 2025",
        "Three new product lines launching in H2",
        "Market expansion into APAC and LATAM regions",
        "Customer retention goal: 94% net revenue retention",
        "Engineering headcount growth: 35 new hires planned",
    ])

    # Slide 3: Market Analysis
    add_content_slide(prs, "Market Analysis", [
        "Total addressable market: $8.2B (growing 18% YoY)",
        "Primary competitors: Nexus Corp, Avelon Systems, TrueGrid",
        "Key differentiator: AI-powered automation reduces onboarding by 60%",
        "Enterprise segment growing fastest at 24% CAGR",
        "SMB churn remains elevated at 3.2% monthly",
    ])

    # Slide 4: Product Roadmap
    add_content_slide(prs, "Product Roadmap", [
        "Q1: Platform reliability improvements (99.95% uptime SLA)",
        "Q2: Analytics Dashboard v3 with predictive insights",
        "Q3: Mobile companion app for iOS and Android",
        "Q4: Enterprise SSO and advanced RBAC controls",
        "Ongoing: API ecosystem expansion and partner integrations",
    ])

    # Slide 5: Revenue Projections
    add_content_slide(prs, "Revenue Projections", [
        "Current ARR: $28.5M (as of January 2025)",
        "Q1 target: $31.2M (+9.5% QoQ)",
        "Q2 target: $34.8M (+11.5% QoQ)",
        "Q3 target: $38.1M (+9.5% QoQ)",
        "Q4 target: $42.0M (+10.2% QoQ)",
    ])

    # Slide 6: Customer Segmentation
    add_content_slide(prs, "Customer Segmentation", [
        "Enterprise (500+ seats): 42 accounts, $16.8M ARR",
        "Mid-Market (50-499 seats): 187 accounts, $8.2M ARR",
        "SMB (1-49 seats): 1,240 accounts, $3.5M ARR",
        "Top 10 accounts represent 28% of total revenue",
        "Expansion revenue: 118% net dollar retention in Enterprise",
    ])

    # Slide 7: Competitive Positioning
    add_content_slide(prs, "Competitive Positioning", [
        "Strength: Best-in-class onboarding experience (NPS 72)",
        "Strength: API-first architecture enables rapid integrations",
        "Weakness: Limited offline capabilities vs. Avelon Systems",
        "Opportunity: AI copilot features ahead of all competitors",
        "Threat: Nexus Corp aggressive pricing in mid-market",
    ])

    # Slide 8: Go-to-Market Process — TITLE ONLY, NO CONTENT
    add_title_only_slide(prs, "Go-to-Market Process")

    # Slide 9: Team Structure
    add_content_slide(prs, "Team Structure", [
        "Engineering: 78 engineers across 12 squads",
        "Product: 8 product managers, 4 designers",
        "Sales: 24 AEs, 6 SDRs, 3 Sales Engineers",
        "Customer Success: 15 CSMs, 5 Support Engineers",
        "Marketing: 12 specialists across demand gen and content",
    ])

    # Slide 10: Key Milestones & Next Steps
    add_content_slide(prs, "Key Milestones & Next Steps", [
        "March 15: Analytics Dashboard v3 beta launch",
        "April 30: APAC region go-live (Singapore data center)",
        "June 1: Mobile app public beta",
        "August 15: Enterprise SSO general availability",
        "October 1: Annual customer summit (San Francisco)",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
