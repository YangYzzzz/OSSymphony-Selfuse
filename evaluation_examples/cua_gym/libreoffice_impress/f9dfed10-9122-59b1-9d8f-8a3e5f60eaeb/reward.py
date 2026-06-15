"""
Reward Script: Create chevron process flow on slide 8
Task ID: impress_exec_054
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): 5 chevron-type shapes exist on slide 8
  Component 2 (0.30): Correct labels in order (Ideation, Validation, Development, Launch, Scale)
  Component 3 (0.20): Shapes evenly spaced horizontally
  Component 4 (0.20): Progressive blue gradient from light (#B3D9FF) to dark (#003366)
"""

import os

from pptx import Presentation
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_054'

EXPECTED_LABELS = ['Ideation', 'Validation', 'Development', 'Launch', 'Scale']
EXPECTED_FIRST_COLOR = 'B3D9FF'
EXPECTED_LAST_COLOR = '003366'


def get_chevron_shapes(slide):
    """Return list of shapes that are chevron/arrow preset geometry on the given slide."""
    chevrons = []
    for shape in slide.shapes:
        el = shape._element
        prstGeom = el.find('.//' + qn('a:prstGeom'))
        if prstGeom is not None:
            prst = prstGeom.get('prst', '')
            # Accept chevron, rightArrow, homePlate, notchedRightArrow, stripedRightArrow, pentagon
            if prst in ('chevron', 'rightArrow', 'homePlate', 'notchedRightArrow',
                        'stripedRightArrow', 'pentagon'):
                chevrons.append((shape, prst))
    return chevrons


def get_shape_fill_color(shape):
    """Extract the primary solid fill color (srgbClr) from a shape's spPr."""
    el = shape._element
    # Look for spPr > solidFill > srgbClr
    spPr = el.find(qn('p:spPr'))
    if spPr is None:
        # Try drawingML namespace
        spPr = el.find('.//' + qn('a:solidFill'))
        if spPr is not None:
            srgb = spPr.find(qn('a:srgbClr'))
            if srgb is not None:
                return srgb.get('val', '').upper()
        return None

    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is not None:
        srgb = solidFill.find(qn('a:srgbClr'))
        if srgb is not None:
            return srgb.get('val', '').upper()
    return None


def color_hex_to_brightness(hex_color):
    """Convert hex color to perceived brightness (0=dark, 255=bright)."""
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return 0.299 * r + 0.587 * g + 0.114 * b


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 8 slides
    if len(prs.slides) < 8:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 8")
        print("REWARD: 0.0")
        return 0.0

    slide8 = prs.slides[7]  # 0-indexed

    # Component 1: 5 chevron/arrow shapes exist on slide 8 (0.30 points)
    try:
        chevron_shapes = get_chevron_shapes(slide8)
        num_chevrons = len(chevron_shapes)
        if num_chevrons == 5:
            print(f"PASS: Component 1 -- Found 5 chevron shapes on slide 8 (0.30 pts)")
            total_score += 0.30
        elif num_chevrons >= 3:
            partial = 0.15
            print(f"PARTIAL: Component 1 -- Found {num_chevrons} chevron shapes, expected 5 ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- Found {num_chevrons} chevron shapes on slide 8, expected 5")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # If no chevrons found, cannot proceed with label/color checks
    if len(chevron_shapes) == 0:
        print("No chevron shapes found, skipping remaining checks")
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Sort chevrons by left position (left to right)
    chevron_shapes.sort(key=lambda x: x[0].left)

    # Component 2: Correct labels in left-to-right order (0.30 points)
    try:
        actual_labels = []
        for shape, _ in chevron_shapes:
            text = ''
            if hasattr(shape, 'text_frame'):
                text = shape.text_frame.text.strip()
            elif hasattr(shape, 'text'):
                text = shape.text.strip()
            actual_labels.append(text)

        matched = 0
        for i, expected in enumerate(EXPECTED_LABELS):
            if i < len(actual_labels):
                if actual_labels[i].lower() == expected.lower():
                    matched += 1

        if matched == 5:
            print(f"PASS: Component 2 -- All 5 labels correct in order: {actual_labels} (0.30 pts)")
            total_score += 0.30
        elif matched >= 3:
            partial = round(0.30 * matched / 5, 2)
            print(f"PARTIAL: Component 2 -- {matched}/5 labels matched. Actual: {actual_labels} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 -- Labels: {actual_labels}, expected: {EXPECTED_LABELS}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Shapes evenly spaced horizontally (0.20 points)
    try:
        if len(chevron_shapes) >= 3:
            lefts = [s.left for s, _ in chevron_shapes]
            gaps = [lefts[i+1] - lefts[i] for i in range(len(lefts)-1)]
            if len(gaps) > 0:
                avg_gap = sum(gaps) / len(gaps)
                max_deviation = max(abs(g - avg_gap) for g in gaps)
                # Allow 10% tolerance on gap uniformity
                tolerance = avg_gap * 0.10
                if max_deviation <= tolerance:
                    print(f"PASS: Component 3 -- Shapes evenly spaced, avg gap={avg_gap:.0f} EMU, max deviation={max_deviation:.0f} (0.20 pts)")
                    total_score += 0.20
                else:
                    print(f"FAIL: Component 3 -- Uneven spacing. Gaps: {gaps}, avg={avg_gap:.0f}, max_dev={max_deviation:.0f}, tolerance={tolerance:.0f}")
            else:
                print(f"FAIL: Component 3 -- Not enough shapes for gap calculation")
        else:
            print(f"FAIL: Component 3 -- Need at least 3 chevron shapes for spacing check")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Progressive blue gradient from light to dark (0.20 points)
    try:
        colors = []
        for shape, _ in chevron_shapes:
            color = get_shape_fill_color(shape)
            colors.append(color)

        print(f"  DEBUG: Detected fill colors (L-to-R): {colors}")

        valid_colors = [c for c in colors if c is not None and len(c) == 6]
        if len(valid_colors) >= 2:
            # Check that brightness decreases left to right (light -> dark)
            brightnesses = [color_hex_to_brightness(c) for c in valid_colors]
            is_decreasing = all(brightnesses[i] >= brightnesses[i+1] for i in range(len(brightnesses)-1))

            # Check first color close to light blue and last close to dark blue
            first_color = valid_colors[0]
            last_color = valid_colors[-1]

            first_bright = color_hex_to_brightness(first_color)
            last_bright = color_hex_to_brightness(last_color)

            # First should be light (brightness > 150), last should be dark (brightness < 100)
            first_is_light = first_bright > 140
            last_is_dark = last_bright < 100

            points = 0.0
            if is_decreasing:
                points += 0.10
                print(f"  SUB-PASS: Colors decrease in brightness L-to-R: {[f'{b:.0f}' for b in brightnesses]}")
            else:
                print(f"  SUB-FAIL: Colors do NOT decrease in brightness: {[f'{b:.0f}' for b in brightnesses]}")

            if first_is_light and last_is_dark:
                points += 0.10
                print(f"  SUB-PASS: First is light ({first_color}, bright={first_bright:.0f}), last is dark ({last_color}, bright={last_bright:.0f})")
            else:
                print(f"  SUB-FAIL: First light={first_is_light} ({first_color}, bright={first_bright:.0f}), last dark={last_is_dark} ({last_color}, bright={last_bright:.0f})")

            if points > 0:
                print(f"PASS: Component 4 -- Progressive gradient verified ({points} pts)")
                total_score += points
            else:
                print(f"FAIL: Component 4 -- No gradient progression detected")
        else:
            print(f"FAIL: Component 4 -- Could not extract enough fill colors: {colors}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
