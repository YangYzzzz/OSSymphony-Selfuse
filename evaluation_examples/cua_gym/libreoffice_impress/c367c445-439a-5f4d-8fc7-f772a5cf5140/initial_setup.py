"""
Initial Setup: Public Health presentation with 9 slides, no transitions.
Task ID: impress_stu_029
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_029'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_text(slide, title_text, subtitle_text=None):
    """Helper to set title and optional subtitle on a slide."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    if subtitle_text and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle_text


def add_body_text(slide, lines):
    """Add bullet-point style text to the content placeholder."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:  # body/content placeholder
            tf = ph.text_frame
            tf.clear()
            for i, line in enumerate(lines):
                if i == 0:
                    tf.paragraphs[0].text = line
                else:
                    p = tf.add_paragraph()
                    p.text = line
            return


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1 - Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Public Health in the 21st Century"
    slide1.placeholders[1].text = "Trends, Challenges, and Opportunities\nDr. Emily Nakamura | March 2026"

    # Slide 2 - Overview
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Presentation Overview"
    add_body_text(slide2, [
        "Global health landscape and emerging trends",
        "Infectious disease preparedness post-COVID",
        "Non-communicable disease burden worldwide",
        "Mental health as a public health priority",
        "Digital health and telemedicine expansion",
        "Health equity and social determinants",
        "Policy recommendations and next steps",
    ])

    # Slide 3 - Global Health Landscape
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Global Health Landscape"
    add_body_text(slide3, [
        "World population reached 8.1 billion in 2025",
        "Average life expectancy: 73.4 years globally",
        "Healthcare spending: $9.8 trillion worldwide",
        "1.8 billion people lack access to basic healthcare",
        "Sub-Saharan Africa faces 24% of disease burden with 1% of health workforce",
        "Climate change projected to cause 250,000 additional deaths annually by 2030",
    ])

    # Slide 4 - Infectious Disease Preparedness
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Infectious Disease Preparedness"
    add_body_text(slide4, [
        "Post-COVID pandemic readiness frameworks established in 147 countries",
        "WHO Pandemic Prevention Treaty negotiations ongoing since 2024",
        "mRNA vaccine platforms enable 100-day response timeline",
        "Antimicrobial resistance: 4.95 million associated deaths in 2023",
        "Global Health Security Index average score: 38.9 out of 100",
        "Surveillance networks expanded with wastewater genomic monitoring",
    ])

    # Slide 5 - Non-Communicable Diseases
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Non-Communicable Diseases (NCDs)"
    add_body_text(slide5, [
        "NCDs responsible for 74% of all deaths globally",
        "Cardiovascular disease: 17.9 million deaths annually",
        "Diabetes prevalence doubled since 2000 to 537 million adults",
        "Cancer incidence rising 1.4% per year in low-income countries",
        "Obesity rates tripled in children and adolescents since 1975",
        "Prevention through lifestyle modification could avert 80% of NCDs",
    ])

    # Slide 6 - Mental Health
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Mental Health as a Priority"
    add_body_text(slide6, [
        "1 in 8 people globally lives with a mental health condition",
        "Depression is the leading cause of disability worldwide",
        "Youth mental health crisis: 25% increase in anxiety disorders since 2020",
        "Treatment gap exceeds 75% in low- and middle-income countries",
        "Economic cost: $2.5 trillion annually in lost productivity",
        "Integration of mental health into primary care shows 40% better outcomes",
    ])

    # Slide 7 - Digital Health
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Digital Health and Telemedicine"
    add_body_text(slide7, [
        "Telemedicine market grew to $115 billion by 2025",
        "AI-assisted diagnostics achieving 94% accuracy in dermatology",
        "Wearable health devices: 1.1 billion users globally",
        "Electronic health records adopted by 89% of hospitals in OECD countries",
        "Remote patient monitoring reduced hospital readmissions by 38%",
        "Digital divide remains: 2.7 billion people still lack internet access",
    ])

    # Slide 8 - Health Equity
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Health Equity and Social Determinants"
    add_body_text(slide8, [
        "Social determinants account for 30-55% of health outcomes",
        "Income inequality correlates with 15-year life expectancy gap",
        "Rural communities face 68% fewer physicians per capita",
        "Maternal mortality 2.9x higher in Black women in the United States",
        "Food insecurity affects 828 million people worldwide",
        "Universal Health Coverage reached by only 47% of countries",
    ])

    # Slide 9 - Thank You
    slide9 = prs.slides.add_slide(prs.slide_layouts[0])
    slide9.shapes.title.text = "Thank You"
    slide9.placeholders[1].text = "Questions & Discussion\nContact: e.nakamura@publichealth.org"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
