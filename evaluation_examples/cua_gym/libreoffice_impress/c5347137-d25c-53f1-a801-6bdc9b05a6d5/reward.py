"""
Reward Script: Stacked bar chart on slide 7 with time allocation data
Task ID: impress_stu_066
Domain: libreoffice_impress
Scoring:
  Component 1: Chart exists on slide 7 and is stacked bar type (0.20)
  Component 2: Chart title is 'Weekly Time Distribution (hours)' (0.10)
  Component 3: Correct 5 categories (Monday-Friday) (0.15)
  Component 4: Correct 4 series names (Classes, Studying, Work, Free Time) (0.15)
  Component 5: Correct data values for all series (0.20)
  Component 6: Series colors approximately match spec (0.10)
  Component 7: Data labels present on chart segments (0.10)
"""

import os
import zipfile
import re
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_066'

# Expected values from task spec
EXPECTED_CATEGORIES = ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday']
EXPECTED_SERIES = {
    'Classes':   [4.0, 3.0, 4.0, 3.0, 2.0],
    'Studying':  [3.0, 4.0, 3.0, 4.0, 2.0],
    'Work':      [2.0, 0.0, 2.0, 0.0, 4.0],
    'Free Time': [3.0, 5.0, 3.0, 5.0, 4.0],
}
# Expected colors: blue, green, orange, gray (approximate RGB ranges)
EXPECTED_COLOR_RANGES = {
    'Classes':   {'r': (0, 120), 'g': (50, 180), 'b': (150, 255)},   # blue-ish
    'Studying':  {'r': (0, 130), 'g': (100, 255), 'b': (0, 130)},    # green-ish
    'Work':      {'r': (180, 255), 'g': (80, 180), 'b': (0, 100)},   # orange-ish
    'Free Time': {'r': (100, 200), 'g': (100, 200), 'b': (100, 200)},# gray-ish
}


def color_in_range(rgb_str, expected_range):
    """Check if a hex RGB string falls within expected range."""
    try:
        r = int(rgb_str[0:2], 16)
        g = int(rgb_str[2:4], 16)
        b = int(rgb_str[4:6], 16)
        return (expected_range['r'][0] <= r <= expected_range['r'][1] and
                expected_range['g'][0] <= g <= expected_range['g'][1] and
                expected_range['b'][0] <= b <= expected_range['b'][1])
    except Exception:
        return False


def check_data_labels_via_xml(file_path):
    """Check if data labels (showVal) are set in chart XML."""
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            chart_files = [f for f in zf.namelist() if 'chart' in f.lower() and f.endswith('.xml')]
            for cf in chart_files:
                with zf.open(cf) as f:
                    content = f.read().decode('utf-8')
                    # Count showVal="1" occurrences (one per series with data labels)
                    show_val_matches = re.findall(r'showVal\s+val="1"', content)
                    if len(show_val_matches) >= 4:
                        return True, len(show_val_matches)
            return False, 0
    except Exception as e:
        return False, str(e)


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Check we have at least 7 slides
    if len(prs.slides) < 7:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 7")
        print("REWARD: 0.0")
        return 0.0

    slide7 = prs.slides[6]  # 0-indexed

    # Find chart shape on slide 7
    chart_shape = None
    for shape in slide7.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            chart_shape = shape
            break

    if chart_shape is None:
        print("FAIL: No chart found on slide 7")
        print("REWARD: 0.0")
        return 0.0

    chart = chart_shape.chart

    # Component 1: Chart is stacked bar type (0.20 points)
    try:
        # COLUMN_STACKED = 52 in python-pptx
        chart_type = chart.chart_type
        if chart_type == XL_CHART_TYPE.COLUMN_STACKED or chart_type == XL_CHART_TYPE.BAR_STACKED:
            print(f"PASS: Component 1 — Chart is stacked type (enum={chart_type}) (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 — Expected stacked bar/column chart, got enum={chart_type}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Chart title (0.10 points)
    try:
        if chart.has_title:
            title_text = chart.chart_title.text_frame.text.strip()
            if title_text == 'Weekly Time Distribution (hours)':
                print(f"PASS: Component 2 — Chart title matches exactly (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 2 — Expected 'Weekly Time Distribution (hours)', found '{title_text}'")
        else:
            print("FAIL: Component 2 — Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Correct 5 categories (0.15 points)
    try:
        plot = chart.plots[0]
        actual_cats = [str(c) for c in plot.categories]
        if actual_cats == EXPECTED_CATEGORIES:
            print(f"PASS: Component 3 — Categories match: {actual_cats} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 3 — Expected {EXPECTED_CATEGORIES}, found {actual_cats}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Correct 4 series names (0.15 points)
    try:
        # Extract series names via XML since python-pptx may not expose them easily
        series_names = []
        with zipfile.ZipFile(file_path, 'r') as zf:
            chart_files = [f for f in zf.namelist() if 'chart' in f.lower() and f.endswith('.xml')]
            for cf in chart_files:
                with zf.open(cf) as f:
                    content = f.read().decode('utf-8')
                    names = re.findall(r'<c:tx>.*?<c:v>([^<]+)</c:v>.*?</c:tx>', content, re.DOTALL)
                    if names:
                        series_names = names

        expected_names = list(EXPECTED_SERIES.keys())
        if series_names == expected_names:
            print(f"PASS: Component 4 — Series names match: {series_names} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 — Expected {expected_names}, found {series_names}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Correct data values (0.20 points)
    try:
        expected_series_list = list(EXPECTED_SERIES.values())
        num_series = len(chart.series)
        if num_series != 4:
            print(f"FAIL: Component 5 — Expected 4 series, found {num_series}")
        else:
            mismatched_series = 0
            for si, series in enumerate(chart.series):
                actual_vals = list(series.values)
                expected_vals = expected_series_list[si]
                if actual_vals != expected_vals:
                    print(f"FAIL: Component 5 — Series {si} values: expected {expected_vals}, found {actual_vals}")
                    mismatched_series += 1
            if mismatched_series == 0:
                print(f"PASS: Component 5 — All series data values match (0.20 pts)")
                total_score += 0.20
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Series colors approximately correct (0.10 points)
    try:
        expected_names_list = list(EXPECTED_SERIES.keys())
        colors_correct = 0
        for si, series in enumerate(chart.series):
            try:
                fill = series.format.fill
                if fill.type is not None:
                    rgb_str = str(fill.fore_color.rgb)
                    name = expected_names_list[si] if si < len(expected_names_list) else f"Series{si}"
                    if name in EXPECTED_COLOR_RANGES:
                        if color_in_range(rgb_str, EXPECTED_COLOR_RANGES[name]):
                            colors_correct += 1
                            print(f"  Series '{name}' color {rgb_str} — in range")
                        else:
                            print(f"  Series '{name}' color {rgb_str} — out of range")
                    else:
                        print(f"  Series '{name}' — no expected color range")
                else:
                    print(f"  Series {si} — no solid fill")
            except Exception as e2:
                print(f"  Series {si} color check error: {e2}")

        if colors_correct >= 3:
            print(f"PASS: Component 6 — {colors_correct}/4 series colors correct (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 6 — Only {colors_correct}/4 series colors in range")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Data labels present (0.10 points)
    try:
        has_labels, info = check_data_labels_via_xml(file_path)
        if has_labels:
            print(f"PASS: Component 7 — Data labels present (showVal count: {info}) (0.10 pts)")
            total_score += 0.10
        else:
            # Also check via python-pptx plot level
            plot = chart.plots[0]
            if plot.has_data_labels and plot.data_labels.show_value:
                print(f"PASS: Component 7 — Data labels present at plot level (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 7 — No data labels found (info: {info})")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'

persist_app_state()

if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
