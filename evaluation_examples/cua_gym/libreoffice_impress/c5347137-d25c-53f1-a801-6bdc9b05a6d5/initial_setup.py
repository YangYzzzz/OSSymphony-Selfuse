"""
Initial Setup: Student Life Survey presentation with 9 slides.
Slide 7 titled 'Time Management Analysis' is empty (no chart).
Task ID: impress_stu_066
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_066'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
    return slide


def add_title_only_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(
        prs,
        "Student Life Survey 2025",
        "Annual Report on Campus Living & Academic Experience"
    )

    # Slide 2: Survey Overview
    add_content_slide(prs, "Survey Overview", [
        "Total Respondents: 2,847 students across 6 faculties",
        "Survey Period: September - November 2025",
        "Response Rate: 68.3% of enrolled students",
        "Methodology: Online questionnaire with follow-up interviews",
        "Key Focus Areas: Academics, Housing, Social Life, Finances, Well-being",
    ])

    # Slide 3: Demographics
    add_content_slide(prs, "Respondent Demographics", [
        "Freshmen (Year 1): 32% — 912 students",
        "Sophomores (Year 2): 26% — 740 students",
        "Juniors (Year 3): 23% — 655 students",
        "Seniors (Year 4): 19% — 540 students",
        "Gender Split: 54% Female, 43% Male, 3% Non-binary",
        "Residential: 61% On-Campus, 39% Off-Campus",
    ])

    # Slide 4: Academic Satisfaction
    add_content_slide(prs, "Academic Satisfaction Ratings", [
        "Overall Academic Experience: 4.1 / 5.0",
        "Quality of Teaching: 3.9 / 5.0",
        "Library & Research Resources: 4.3 / 5.0",
        "Course Availability: 3.5 / 5.0",
        "Academic Advising: 3.7 / 5.0",
        "Career Preparation Support: 3.2 / 5.0",
    ])

    # Slide 5: Financial Overview
    add_content_slide(prs, "Student Financial Snapshot", [
        "Average Monthly Budget: $1,450",
        "Tuition as % of Family Income: 28% (median)",
        "Students with Part-Time Jobs: 47%",
        "Average Weekly Work Hours: 14.6 hours",
        "Students Receiving Financial Aid: 63%",
        "Average Student Loan Debt (Seniors): $27,300",
    ])

    # Slide 6: Campus Life Highlights
    add_content_slide(prs, "Campus Life Highlights", [
        "Active in Student Organizations: 72%",
        "Average Club Memberships: 2.3 per student",
        "Regular Gym / Recreation Use: 58%",
        "Attended Campus Events (per semester): 8.4",
        "Satisfaction with Dining Services: 3.6 / 5.0",
        "Satisfaction with Mental Health Resources: 3.1 / 5.0",
    ])

    # Slide 7: Time Management Analysis — EMPTY (title only, NO chart)
    add_title_only_slide(prs, "Time Management Analysis")

    # Slide 8: Housing & Transportation
    add_content_slide(prs, "Housing & Transportation", [
        "Average Monthly Rent (Off-Campus): $785",
        "Commute Time (Off-Campus avg): 24 minutes",
        "Primary Transport: Walking (41%), Bus (29%), Car (18%), Bike (12%)",
        "Satisfaction with Dorm Conditions: 3.4 / 5.0",
        "Students with Roommates: 78%",
    ])

    # Slide 9: Key Takeaways & Recommendations
    add_content_slide(prs, "Key Takeaways & Recommendations", [
        "Expand mental health services — lowest satisfaction score area",
        "Increase course section availability for popular majors",
        "Invest in career services and internship programs",
        "Improve dining options based on dietary preference data",
        "Provide additional financial literacy workshops",
        "Continue fostering active campus community engagement",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
