"""
Reward Script: Verify resource links slide with 6 clickable hyperlink cards
Task ID: impress_ps_024
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): 6 rounded rectangle shapes on slide 8
  Component 2 (0.4): Correct hyperlinks on all 6 resource cards
  Component 3 (0.2): Text formatting (underlined, blue, 16pt)
  Component 4 (0.1): Card appearance (fill + shadow)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_024'

# Expected resource cards: name -> URL
EXPECTED_RESOURCES = {
    'HR Portal': 'https://hr.example.com',
    'IT Help Desk': 'https://helpdesk.example.com',
    'Benefits Guide': 'https://benefits.example.com',
    'Training Platform': 'https://learn.example.com',
    'Expense System': 'https://expense.example.com',
    'Company Wiki': 'https://wiki.example.com',
}


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: File must have at least 8 slides
    if len(prs.slides) < 8:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 8")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[7]  # 0-indexed, slide 8

    # Collect all AUTO_SHAPE shapes (the resource cards) - exclude pre-existing shapes
    auto_shapes = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            auto_shapes.append(shape)

    # Component 1: 6 rounded rectangle shapes exist on slide 8 (0.3 points)
    try:
        num_cards = len(auto_shapes)
        if num_cards >= 6:
            print(f"PASS: Component 1 - Found {num_cards} auto shapes on slide 8 (0.3 pts)")
            total_score += 0.3
        elif num_cards >= 3:
            partial = 0.15
            print(f"PARTIAL: Component 1 - Found {num_cards}/6 auto shapes (0.15 pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 - Found {num_cards}/6 auto shapes on slide 8")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Correct hyperlinks on all 6 resource cards (0.4 points)
    # Each correct name+URL pair = 0.4/6 points
    try:
        found_links = {}
        for shape in auto_shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        text = (run.text or "").strip()
                        hlnk = run.hyperlink
                        url = hlnk.address if hlnk and hlnk.address else None
                        if text and url:
                            found_links[text] = url

        link_score = 0.0
        per_link = 0.4 / 6.0
        matched = 0
        for name, expected_url in EXPECTED_RESOURCES.items():
            if name in found_links:
                actual_url = found_links[name].rstrip('/')
                expected_clean = expected_url.rstrip('/')
                if actual_url == expected_clean:
                    matched += 1
                    link_score += per_link
                    print(f"  PASS: '{name}' -> {actual_url}")
                else:
                    print(f"  FAIL: '{name}' URL mismatch: expected {expected_clean}, got {actual_url}")
            else:
                print(f"  FAIL: '{name}' not found in shape hyperlinks")

        link_score = round(link_score, 4)
        if matched == 6:
            print(f"PASS: Component 2 - All 6 hyperlinks correct ({link_score} pts)")
        elif matched > 0:
            print(f"PARTIAL: Component 2 - {matched}/6 hyperlinks correct ({link_score} pts)")
        else:
            print(f"FAIL: Component 2 - 0/6 hyperlinks correct")
        if link_score > 0:
            total_score += link_score
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Text formatting - underlined, blue color, ~16pt (0.2 points)
    # Check across all auto shapes
    try:
        format_ok_count = 0
        for shape in auto_shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not (run.text or "").strip():
                        continue
                    is_underlined = (run.font.underline is True)
                    is_blue = False
                    try:
                        if run.font.color.type is not None:
                            rgb_str = str(run.font.color.rgb).upper()
                            # Accept various blue shades
                            r = int(rgb_str[0:2], 16)
                            g = int(rgb_str[2:4], 16)
                            b = int(rgb_str[4:6], 16)
                            is_blue = (b > 150 and b > r and b > g)
                    except Exception:
                        pass
                    has_size = False
                    if run.font.size is not None:
                        pt_size = run.font.size / 12700
                        has_size = (14 <= pt_size <= 18)

                    if is_underlined and is_blue and has_size:
                        format_ok_count += 1
                        break  # one run per shape is enough

        format_score = min(format_ok_count / 6.0, 1.0) * 0.2
        format_score = round(format_score, 4)
        if format_ok_count >= 6:
            print(f"PASS: Component 3 - {format_ok_count}/6 shapes have correct formatting ({format_score} pts)")
        elif format_ok_count > 0:
            print(f"PARTIAL: Component 3 - {format_ok_count}/6 shapes formatted correctly ({format_score} pts)")
        else:
            print(f"FAIL: Component 3 - 0/6 shapes formatted correctly")
        if format_score > 0:
            total_score += format_score
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Card appearance - fill and shadow (0.1 points)
    try:
        cards_with_fill = 0
        cards_with_shadow = 0
        for shape in auto_shapes:
            # Check fill
            try:
                fill = shape.fill
                if fill.type is not None and fill.type == 1:  # SOLID
                    cards_with_fill += 1
            except Exception:
                pass
            # Check shadow via XML
            try:
                from lxml import etree
                ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                shadow = shape._element.find('.//a:effectLst/a:outerShdw', ns)
                if shadow is not None:
                    cards_with_shadow += 1
            except Exception:
                pass

        appearance_score = 0.0
        if cards_with_fill >= 5:
            appearance_score += 0.05
            print(f"  PASS: {cards_with_fill}/6 cards have fill color")
        else:
            print(f"  FAIL: Only {cards_with_fill}/6 cards have fill color")
        if cards_with_shadow >= 5:
            appearance_score += 0.05
            print(f"  PASS: {cards_with_shadow}/6 cards have shadow effect")
        else:
            print(f"  FAIL: Only {cards_with_shadow}/6 cards have shadow effect")

        if appearance_score > 0:
            print(f"PASS: Component 4 - Card appearance ({appearance_score} pts)")
            total_score += appearance_score
        else:
            print(f"FAIL: Component 4 - Card appearance missing")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
