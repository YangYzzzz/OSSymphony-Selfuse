"""
Initial Setup: New Hire Onboarding presentation with 10 slides.
Slide 8 has title 'Useful Resources' but empty content area.
Task ID: impress_ps_024
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_024'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.paragraphs[0].text = point
        else:
            p = tf.add_paragraph()
            p.text = point
    return slide


def add_title_only_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(prs, "Welcome to Meridian Technologies",
                    "New Hire Onboarding Program — Q2 2025")

    # Slide 2: Agenda
    add_content_slide(prs, "Onboarding Agenda", [
        "Company Overview & Mission",
        "Team Structure & Key Contacts",
        "IT Setup & Access Credentials",
        "Benefits Enrollment & Payroll",
        "First Week Schedule",
        "Useful Resources",
        "Q&A Session",
    ])

    # Slide 3: Company Overview
    add_content_slide(prs, "About Meridian Technologies", [
        "Founded in 2012, headquartered in Austin, TX",
        "Over 3,200 employees across 14 global offices",
        "Specializing in enterprise cloud solutions and data analytics",
        "Named a Great Place to Work for 5 consecutive years",
        "Core values: Innovation, Integrity, Inclusion, Impact",
    ])

    # Slide 4: Team Structure
    add_content_slide(prs, "Your Team & Key Contacts", [
        "Direct Manager: Rachel Torres (rachel.torres@meridian.com)",
        "HR Business Partner: James Liu (james.liu@meridian.com)",
        "IT Support Lead: Priya Sharma (priya.sharma@meridian.com)",
        "Buddy/Mentor: Alex Novak (alex.novak@meridian.com)",
        "Office Manager: Diana Reeves (diana.reeves@meridian.com)",
    ])

    # Slide 5: IT Setup
    add_content_slide(prs, "IT Setup & Access", [
        "Laptop: Pre-configured MacBook Pro or Dell XPS (your choice)",
        "Email: Outlook Web — login at mail.meridian.com",
        "VPN: Cisco AnyConnect — config guide in IT Wiki",
        "Slack workspace: meridian-tech.slack.com",
        "Code repos: GitHub Enterprise — github.meridian.com",
        "Password manager: 1Password — invitation sent to your email",
    ])

    # Slide 6: Benefits
    add_content_slide(prs, "Benefits & Payroll", [
        "Health, dental, and vision insurance (effective Day 1)",
        "401(k) with 4% company match (eligible after 30 days)",
        "20 PTO days + 12 company holidays per year",
        "Annual learning stipend of $2,500",
        "Gym membership reimbursement up to $75/month",
        "Payroll processed biweekly — direct deposit setup in HR Portal",
    ])

    # Slide 7: First Week Schedule
    add_content_slide(prs, "Your First Week", [
        "Monday: Orientation session (9 AM), badge pickup, workspace tour",
        "Tuesday: IT setup, Slack intro, meet your team",
        "Wednesday: Product deep-dive with Engineering leads",
        "Thursday: Benefits enrollment workshop, 1:1 with manager",
        "Friday: Team lunch, goal-setting session, first-week retrospective",
    ])

    # Slide 8: Useful Resources — title only, empty content area
    # This is the slide the agent needs to modify
    slide8 = add_title_only_slide(prs, "Useful Resources")
    # Intentionally left empty — no shapes, textboxes, or hyperlinks below the title

    # Slide 9: FAQ
    add_content_slide(prs, "Frequently Asked Questions", [
        "Where do I park? — Garage B, levels 2-4 (badge access required)",
        "What's the dress code? — Business casual; Fridays are casual",
        "How do I book a meeting room? — Use the Room Finder in Outlook",
        "Who do I contact for facility issues? — facilities@meridian.com",
        "Can I work remotely? — Hybrid schedule (3 office / 2 remote days)",
    ])

    # Slide 10: Q&A / Closing
    add_title_slide(prs, "Questions & Next Steps",
                    "We're excited to have you on the team!\nReach out anytime: onboarding@meridian.com")

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
