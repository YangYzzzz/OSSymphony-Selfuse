"""
Reward Script: Scatter plot on slide 5 with specific data points, title, and axis labels
Task ID: impress_stu_046
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Scatter chart exists on slide 5
  Component 2 (0.25): Chart title is 'Study Hours vs Exam Performance'
  Component 3 (0.20): Axis labels: X='Hours Studied', Y='Exam Score (%)'
  Component 4 (0.30): 10 data points match specified (X,Y) pairs
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_046'

# Expected data points from the task instruction
EXPECTED_DATA = [
    (2, 55), (4, 65), (5, 70), (6, 75), (7, 80),
    (8, 85), (9, 88), (10, 92), (3, 58), (6, 72)
]

def persist_app_state(domain):
    """Save any unsaved GUI edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
    except ImportError:
        print("CRITICAL: python-pptx not installed")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide5 = prs.slides[4]  # 0-indexed

    # Find chart shape on slide 5
    chart_shape = None
    for shape in slide5.shapes:
        if shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Scatter chart exists on slide 5 (0.25 points)
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            # Check it's a scatter chart type (XY_SCATTER variants)
            chart_type_val = chart.chart_type
            # XY_SCATTER = -4169, XY_SCATTER_LINES = 74, XY_SCATTER_LINES_NO_MARKERS = 75,
            # XY_SCATTER_SMOOTH = 72, XY_SCATTER_SMOOTH_NO_MARKERS = 73
            scatter_types = {-4169, 72, 73, 74, 75}
            if chart_type_val in scatter_types:
                print(f"PASS: Component 1 — Scatter chart found on slide 5, type={chart_type_val} (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 — Chart on slide 5 is not a scatter type, got type={chart_type_val}")
        else:
            print("FAIL: Component 1 — No chart found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # If no chart found, remaining checks will fail; skip them gracefully
    if chart_shape is None:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    chart = chart_shape.chart

    # Component 2: Chart title is 'Study Hours vs Exam Performance' (0.25 points)
    try:
        if chart.has_title:
            title_text = chart.chart_title.text_frame.text.strip()
            if title_text == "Study Hours vs Exam Performance":
                print(f"PASS: Component 2 — Chart title matches exactly (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 2 — Expected title 'Study Hours vs Exam Performance', found '{title_text}'")
        else:
            print("FAIL: Component 2 — Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Axis labels correct (0.20 points)
    # X-axis = 'Hours Studied' (0.10), Y-axis = 'Exam Score (%)' (0.10)
    try:
        axis_score = 0.0

        # X-axis (category axis for scatter)
        try:
            cat_axis = chart.category_axis
            if cat_axis.has_title:
                x_title = cat_axis.axis_title.text_frame.text.strip()
                if x_title == "Hours Studied":
                    print(f"PASS: Component 3a — X-axis label 'Hours Studied' correct (0.10 pts)")
                    axis_score += 0.10
                else:
                    print(f"FAIL: Component 3a — X-axis label expected 'Hours Studied', found '{x_title}'")
            else:
                print("FAIL: Component 3a — X-axis has no title")
        except Exception as e:
            print(f"ERROR: Component 3a — {e}")

        # Y-axis (value axis for scatter)
        try:
            val_axis = chart.value_axis
            if val_axis.has_title:
                y_title = val_axis.axis_title.text_frame.text.strip()
                if y_title == "Exam Score (%)":
                    print(f"PASS: Component 3b — Y-axis label 'Exam Score (%)' correct (0.10 pts)")
                    axis_score += 0.10
                else:
                    print(f"FAIL: Component 3b — Y-axis label expected 'Exam Score (%)', found '{y_title}'")
            else:
                print("FAIL: Component 3b — Y-axis has no title")
        except Exception as e:
            print(f"ERROR: Component 3b — {e}")

        if axis_score > 0:
            total_score += axis_score
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Data points match (0.30 points)
    # Extract X and Y values from the scatter chart series via XML
    try:
        ns = {'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'}
        plot = chart.plots[0]

        if len(plot.series) < 1:
            print("FAIL: Component 4 — No data series found in chart")
        else:
            series = plot.series[0]
            ser_element = series._element

            # Extract X values
            xval_elem = ser_element.find('.//c:xVal', ns)
            yval_elem = ser_element.find('.//c:yVal', ns)

            xvals = []
            yvals = []

            if xval_elem is not None:
                for pt in xval_elem.findall('.//c:pt', ns):
                    v = pt.find('c:v', ns)
                    if v is not None:
                        xvals.append(float(v.text))

            if yval_elem is not None:
                for pt in yval_elem.findall('.//c:pt', ns):
                    v = pt.find('c:v', ns)
                    if v is not None:
                        yvals.append(float(v.text))

            actual_pairs = list(zip(xvals, yvals))
            expected_set = set((float(x), float(y)) for x, y in EXPECTED_DATA)
            actual_set = set(actual_pairs)

            print(f"  Expected {len(EXPECTED_DATA)} data points, found {len(actual_pairs)}")
            print(f"  Actual pairs: {actual_pairs}")

            if len(actual_pairs) == len(EXPECTED_DATA) and actual_set == expected_set:
                print(f"PASS: Component 4 — All 10 data points match exactly (0.30 pts)")
                total_score += 0.30
            elif len(actual_set & expected_set) > 0:
                # Partial credit: fraction of matching points
                match_count = len(actual_set & expected_set)
                partial = 0.30 * (match_count / len(EXPECTED_DATA))
                print(f"PARTIAL: Component 4 — {match_count}/{len(EXPECTED_DATA)} data points match ({partial:.2f} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 4 — No data points match")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
