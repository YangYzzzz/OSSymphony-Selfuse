"""
Initial Setup: Create a 7-slide Stats Analysis presentation with slide 5 empty for chart creation task.
Task ID: impress_stu_046
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_046'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_placeholder(shape, text, font_size=Pt(18), bold=False, color=None):
    """Helper to set text in a placeholder shape."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    for run in p.runs:
        run.font.size = font_size
        run.font.bold = bold
        if color:
            run.font.color.rgb = color


def add_body_text(slide, lines, left=Inches(0.8), top=Inches(1.8), width=Inches(8.4), height=Inches(4.5)):
    """Add a text box with multiple lines/paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Statistical Analysis Report"
    slide1.placeholders[1].text = "Q1 2025 Academic Performance Review\nPrepared by: Dr. Elena Vasquez"

    # --- Slide 2: Introduction ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txTitle = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(1.0))
    tf = txTitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Introduction"
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    add_body_text(slide2, [
        "This report examines the relationship between various study habits",
        "and academic performance among 150 undergraduate students at",
        "Westfield University during the Spring 2025 semester.",
        "",
        "Key research questions:",
        "  - How do study hours correlate with exam scores?",
        "  - What factors predict high academic achievement?",
        "  - Are there diminishing returns to increased study time?",
    ])

    # --- Slide 3: Data Collection Methodology ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txTitle3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(1.0))
    tf3 = txTitle3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Data Collection Methodology"
    p3.alignment = PP_ALIGN.LEFT
    for run in p3.runs:
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    add_body_text(slide3, [
        "Survey Period: January 15 - March 28, 2025",
        "Sample Size: 150 students across 6 departments",
        "Variables Tracked:",
        "  - Daily study hours (self-reported via mobile app)",
        "  - Midterm exam scores (standardized 0-100 scale)",
        "  - Sleep hours, class attendance, prior GPA",
        "",
        "Data cleaned: 12 incomplete records removed (n=138 final)",
    ])

    # --- Slide 4: Summary Statistics ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txTitle4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(1.0))
    tf4 = txTitle4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Summary Statistics"
    p4.alignment = PP_ALIGN.LEFT
    for run in p4.runs:
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Add a small table
    table_shape = slide4.shapes.add_table(5, 3, Inches(1.5), Inches(2.0), Inches(7.0), Inches(3.0))
    table = table_shape.table
    headers = ["Metric", "Study Hours", "Exam Score"]
    data_rows = [
        ["Mean", "6.2", "74.8"],
        ["Median", "6.0", "75.0"],
        ["Std Dev", "2.4", "11.3"],
        ["Range", "1-12", "38-98"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row_data in enumerate(data_rows, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(14)

    # --- Slide 5: Correlation Analysis (EMPTY - task target) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    txTitle5 = slide5.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(1.0))
    tf5 = txTitle5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Correlation Analysis"
    p5.alignment = PP_ALIGN.LEFT
    for run in p5.runs:
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
    # NO chart here - this is the task for the agent

    # --- Slide 6: Key Findings ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    txTitle6 = slide6.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(1.0))
    tf6 = txTitle6.text_frame
    p6 = tf6.paragraphs[0]
    p6.text = "Key Findings"
    p6.alignment = PP_ALIGN.LEFT
    for run in p6.runs:
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    add_body_text(slide6, [
        "1. Strong positive correlation (r = 0.94) between study hours and scores",
        "2. Students studying 8+ hours averaged 86.5% on exams",
        "3. Diminishing returns observed beyond 10 hours of daily study",
        "4. Attendance and prior GPA also significant predictors",
        "",
        "Pearson correlation coefficient: r = 0.94, p < 0.001",
    ])

    # --- Slide 7: Conclusions & Recommendations ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    txTitle7 = slide7.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(1.0))
    tf7 = txTitle7.text_frame
    p7 = tf7.paragraphs[0]
    p7.text = "Conclusions & Recommendations"
    p7.alignment = PP_ALIGN.LEFT
    for run in p7.runs:
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    add_body_text(slide7, [
        "Recommendations for Academic Success:",
        "  - Target 6-8 hours of focused study per day",
        "  - Maintain consistent study schedules over cramming",
        "  - Combine study with regular class attendance",
        "",
        "Next Steps:",
        "  - Longitudinal study across multiple semesters",
        "  - Investigate study quality vs. quantity metrics",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
