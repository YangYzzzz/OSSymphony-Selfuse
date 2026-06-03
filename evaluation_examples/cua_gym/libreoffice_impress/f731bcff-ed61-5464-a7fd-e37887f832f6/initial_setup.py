"""
Initial Setup: Create presentation with 8 slides and one master slide named 'Default'
Task ID: impress_ma_037
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import copy
import zipfile
import shutil
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_037'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Rename the default master slide to 'Default' and set white background ---
    master = prs.slide_masters[0]
    # Set master name via XML
    sldMaster_elem = master.element
    # The cSld element holds the name
    cSld = sldMaster_elem.find(qn('p:cSld'))
    if cSld is not None:
        cSld.set('name', 'Default')

    # Set white background on the master
    bg = sldMaster_elem.find(qn('p:bg'))
    if bg is None:
        bg = sldMaster_elem.makeelement(qn('p:bg'), {})
        # Insert bg before cSld's next sibling or at start
        cSld_idx = list(sldMaster_elem).index(cSld)
        sldMaster_elem.insert(cSld_idx, bg)
    # Clear existing bg content
    for child in list(bg):
        bg.remove(child)
    bgPr = bg.makeelement(qn('p:bgPr'), {})
    solidFill = bgPr.makeelement(qn('a:solidFill'), {})
    srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': 'FFFFFF'})
    solidFill.append(srgbClr)
    bgPr.append(solidFill)
    effectLst = bgPr.makeelement(qn('a:effectLst'), {})
    bgPr.append(effectLst)
    bg.append(bgPr)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q2 2025 Marketing Strategy"
    slide1.placeholders[1].text = "Prepared by Sarah Chen | Marketing Director"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Market Overview & Competitive Analysis"
    p = body2.add_paragraph()
    p.text = "Campaign Performance Metrics"
    p = body2.add_paragraph()
    p.text = "Digital Presence & Social Media Strategy"
    p = body2.add_paragraph()
    p.text = "Budget Allocation & ROI Projections"
    p = body2.add_paragraph()
    p.text = "Team Initiatives & Timeline"

    # --- Slide 3: Market Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Market Overview"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Total addressable market grew 12.3% YoY to $4.8B"
    p = body3.add_paragraph()
    p.text = "Primary competitor Apex Corp holds 28% market share"
    p = body3.add_paragraph()
    p.text = "Our current share: 15.7% (up from 13.2% in Q1)"
    p = body3.add_paragraph()
    p.text = "Key opportunity: enterprise segment expanding at 18% CAGR"

    # --- Slide 4: Campaign Performance ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.7))
    slide4.shapes[-1].text_frame.paragraphs[0].text = "Campaign Performance - Q1 Results"
    run = slide4.shapes[-1].text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Add a table for campaign data
    table_shape = slide4.shapes.add_table(
        5, 4, Inches(0.5), Inches(1.2), Inches(9), Inches(3)
    )
    table = table_shape.table
    headers = ["Campaign", "Impressions", "CTR", "Conversions"]
    data = [
        ["Spring Launch", "2,450,000", "3.2%", "12,840"],
        ["Brand Awareness", "5,120,000", "1.8%", "8,216"],
        ["Product Demo Series", "890,000", "5.7%", "6,423"],
        ["Email Nurture Flow", "340,000", "22.1%", "4,930"],
    ]
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
        for run in table.cell(0, c).text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(12)
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 5: Digital Strategy ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Digital Strategy Roadmap"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Phase 1 (Apr-May): SEO overhaul & content calendar launch"
    p = body5.add_paragraph()
    p.text = "Phase 2 (Jun-Jul): Paid media scaling on LinkedIn & Google"
    p = body5.add_paragraph()
    p.text = "Phase 3 (Aug-Sep): Influencer partnership program rollout"
    p = body5.add_paragraph()
    p.text = "Target: 40% increase in organic traffic by end of Q3"

    # --- Slide 6: Budget Overview ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.7))
    slide6.shapes[-1].text_frame.paragraphs[0].text = "Budget Allocation - Q2 2025"
    run = slide6.shapes[-1].text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(28)
    run.font.bold = True

    tb = slide6.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
    tf = tb.text_frame
    tf.word_wrap = True
    items = [
        "Paid Advertising: $185,000 (37%)",
        "Content Creation: $95,000 (19%)",
        "Events & Sponsorships: $80,000 (16%)",
        "Technology & Tools: $65,000 (13%)",
        "Team Training & Development: $45,000 (9%)",
        "Contingency: $30,000 (6%)",
    ]
    tf.paragraphs[0].text = items[0]
    for item in items[1:]:
        p = tf.add_paragraph()
        p.text = item

    # --- Slide 7: Team Initiatives ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Key Team Initiatives"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "Marcus Johnson: Lead social media rebrand (Apr 15 deadline)"
    p = body7.add_paragraph()
    p.text = "Priya Patel: Website UX audit with external agency"
    p = body7.add_paragraph()
    p.text = "David Kim: CRM integration with marketing automation"
    p = body7.add_paragraph()
    p.text = "Elena Rodriguez: Customer journey mapping workshop series"

    # --- Slide 8: Next Steps ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Next Steps & Timeline"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Week 1-2: Finalize Q2 campaign briefs"
    p = body8.add_paragraph()
    p.text = "Week 3-4: Launch Spring Campaign Phase 2"
    p = body8.add_paragraph()
    p.text = "Month 2: Mid-quarter performance review"
    p = body8.add_paragraph()
    p.text = "Month 3: Prepare Q3 strategy proposal"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Verify
    verify = Presentation(OUTPUT)
    print(f'Slides: {len(verify.slides)}')
    print(f'Masters: {len(verify.slide_masters)}')
    for i, m in enumerate(verify.slide_masters):
        cSld = m.element.find(qn('p:cSld'))
        name = cSld.get('name', '(unnamed)') if cSld is not None else '(unnamed)'
        print(f'  Master {i}: name="{name}"')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
