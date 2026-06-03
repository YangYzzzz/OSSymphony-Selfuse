"""
FINAL REWARD SCRIPT - SUCCESS
Task: I’m freshening up an existing Impress file and want slides 2 through 7 to share the same look. How do I change just those slides so their background is the solid fill color #D5E8D4 — the one LibreOffice labels “Light Green 1” — without touching any of the other slides?
Generated: 2025-09-10 15:43:56
Status: success
Model: azure-o3
Total Steps: 5
"""

import os
import glob
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.dml.color import RGBColor


def verify_backgrounds(file_path: str) -> float:
    """Verify that only slides 2-7 have the Light Green 1 (#D5E8D4) solid background.

    Scoring (progressive):
        • 0.6 points  – Each correctly-formatted slide among slides 2-7 is worth 0.6/6 ≈ 0.10.
        • 0.4 points  – Remaining slides keep their original (non-Light-Green-1) background,
                        each worth an equal share of 0.4.
        The final score is capped at 1.0.
    """
    TARGET_HEX = "D5E8D4"  # LibreOffice “Light Green 1”
    TARGET_RGB = RGBColor(0xD5, 0xE8, 0xD4)

    print(f"Verifying presentation: {file_path}")

    # Load the presentation ---------------------------------------------------
    try:
        prs = Presentation(file_path)
    except Exception as exc:
        print(f"✗ Unable to load presentation: {exc}")
        print("REWARD: 0.0")
        return 0.0  # Cannot evaluate without opening the file

    slide_count = len(prs.slides)
    print(f"Slide count detected: {slide_count}")

    # We need at least 7 slides to perform the required checks
    if slide_count < 7:
        print("✗ Not enough slides present to evaluate task requirements.")
        print("REWARD: 0.0")
        return 0.0

    total_score = 0.0

    # -------------------------------------------------------------------------
    # 1) Verify slides 2-7 (indices 1-6) have correct solid fill               
    # -------------------------------------------------------------------------
    correct_core = 0  # number of target slides with the correct background
    for idx in range(1, 7):  # 0-based indexing → slides 2-7
        slide_number = idx + 1  # for human-readable output
        slide = prs.slides[idx]
        fill = slide.background.fill

        if (
            fill.type == MSO_FILL.SOLID
            and fill.fore_color.rgb is not None
            and str(fill.fore_color.rgb).upper() == TARGET_HEX
        ):
            correct_core += 1
            print(f"✓ Slide {slide_number}: correct Light Green 1 solid background detected")
        else:
            print(f"✗ Slide {slide_number}: WRONG background (expected Light Green 1 solid fill)")

    core_score = (correct_core / 6) * 0.6  # each of the 6 slides worth equal share of 0.6
    total_score += core_score
    print(f"Slides 2-7 background score: {core_score:.2f} / 0.60")

    # -------------------------------------------------------------------------
    # 2) Verify ALL OTHER slides were NOT changed to that specific color       
    # -------------------------------------------------------------------------
    other_indices = list(range(0, 1)) + list(range(7, slide_count))  # slides 1 and 8+ (if any)

    if not other_indices:
        # Edge case: there are no other slides; award full preservation score
        other_score = 0.4
        print("No other slides present; awarding full preservation score (0.40)")
    else:
        others_ok = 0
        for idx in other_indices:
            slide_number = idx + 1
            slide = prs.slides[idx]
            fill = slide.background.fill

            if not (
                fill.type == MSO_FILL.SOLID
                and fill.fore_color.rgb is not None
                and str(fill.fore_color.rgb).upper() == TARGET_HEX
            ):
                others_ok += 1
                print(f"✓ Slide {slide_number}: background correctly NOT set to Light Green 1")
            else:
                print(f"✗ Slide {slide_number}: background INCORRECTLY set to Light Green 1")

        other_score = (others_ok / len(other_indices)) * 0.4

    total_score += other_score
    print(f"Other slides preservation score: {other_score:.2f} / 0.40")

    # -------------------------------------------------------------------------
    final_score = round(min(total_score, 1.0), 2)  # round for cleanliness
    print(f"Total score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score


# -----------------------------------------------------------------------------
# Locate a PPTX file to grade. The evaluation harness usually places the user-
# modified file in /home/user. We search recursively and grade the first one   
# found.                                                                       
# -----------------------------------------------------------------------------

def _main():
    search_paths = glob.glob("/home/user/**/*.pptx", recursive=True)
    if not search_paths:
        print("✗ No .pptx files found for verification.")
        print("REWARD: 0.0")
        return 0.0

    # Use first matching pptx file (there should typically be only one)
    target_file = search_paths[0]
    return verify_backgrounds(target_file)


if __name__ == "__main__":
    _main()

