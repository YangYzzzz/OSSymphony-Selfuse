"""
Reward Script: Add a vertical color bar on the left edge of the master slide
Task ID: impress_ma_036
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): Rectangle shape on master slide at position (0, 0)
  Component 2 (0.3): Correct dimensions (0.5 inches wide, full slide height)
  Component 3 (0.4): Solid fill with color #6C5CE7
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_036'


def persist_app_state(domain: str):
    """Save any unsaved GUI edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Inches, Emu
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Missing library: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Get slide dimensions for reference
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    print(f"INFO: Slide dimensions: {slide_width} x {slide_height} EMU "
          f"({slide_width / 914400:.1f} x {slide_height / 914400:.1f} inches)")

    # Find the color bar shape on the master slide
    # Look for a non-placeholder rectangle on the first slide master
    color_bar = None
    try:
        master = prs.slide_masters[0]
        for shape in master.shapes:
            # Skip placeholders - we want added shapes only
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                continue
            # Look for a shape at the left edge (left ~= 0)
            if shape.left <= Inches(0.1):
                color_bar = shape
                print(f"INFO: Found candidate shape: name='{shape.name}', "
                      f"type={shape.shape_type}, "
                      f"left={shape.left}, top={shape.top}, "
                      f"width={shape.width}, height={shape.height}")
                break
        if color_bar is None:
            print("FAIL: No non-placeholder shape found at left edge of master slide")
    except Exception as e:
        print(f"ERROR: Could not inspect master slide: {e}")

    # Component 1: Rectangle shape on master slide at position (0, 0) (0.3 points)
    try:
        if color_bar is not None:
            # Check position: left should be 0, top should be 0
            left_ok = color_bar.left <= Inches(0.05)  # small tolerance
            top_ok = color_bar.top <= Inches(0.05)     # small tolerance
            if left_ok and top_ok:
                print(f"PASS: Component 1 -- Shape positioned at left edge "
                      f"(left={color_bar.left}, top={color_bar.top}) (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 1 -- Shape not at (0,0): "
                      f"left={color_bar.left} (expect ~0), top={color_bar.top} (expect ~0)")
        else:
            print("FAIL: Component 1 -- No color bar shape found on master slide")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Correct dimensions - 0.5 inches wide, full slide height (0.3 points)
    try:
        if color_bar is not None:
            expected_width = Inches(0.5)   # 457200 EMU
            expected_height = slide_height  # full slide height

            # Use tolerance of 0.5%
            def approx_equal(a, b, tol=0.005):
                if a == b:
                    return True
                if max(abs(a), abs(b)) == 0:
                    return a == b
                return abs(a - b) / max(abs(a), abs(b)) <= tol

            width_ok = approx_equal(color_bar.width, expected_width)
            height_ok = approx_equal(color_bar.height, expected_height)

            if width_ok and height_ok:
                print(f"PASS: Component 2 -- Dimensions correct: "
                      f"width={color_bar.width} (~{expected_width}), "
                      f"height={color_bar.height} (~{expected_height}) (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 2 -- Dimensions mismatch: "
                      f"width={color_bar.width} (expect {expected_width}, ok={width_ok}), "
                      f"height={color_bar.height} (expect {expected_height}, ok={height_ok})")
        else:
            print("FAIL: Component 2 -- No color bar shape found")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Solid fill with color #6C5CE7 (0.4 points)
    try:
        if color_bar is not None:
            fill = color_bar.fill
            if fill.type == 1:  # SOLID fill
                actual_rgb = str(fill.fore_color.rgb).upper()
                expected_rgb = "6C5CE7"
                if actual_rgb == expected_rgb:
                    print(f"PASS: Component 3 -- Fill color is #{actual_rgb} (0.4 pts)")
                    total_score += 0.4
                else:
                    print(f"FAIL: Component 3 -- Fill color is #{actual_rgb}, "
                          f"expected #{expected_rgb}")
            else:
                print(f"FAIL: Component 3 -- Fill type is {fill.type}, expected SOLID (1)")
        else:
            print("FAIL: Component 3 -- No color bar shape found")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
