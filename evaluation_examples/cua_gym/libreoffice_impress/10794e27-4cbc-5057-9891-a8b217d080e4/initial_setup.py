"""
Initial Setup: Workshop Materials presentation with 11 slides, no decorative master slide elements.
Task ID: impress_ma_036
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_036'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard widescreen dimensions
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Advanced Data Analytics Workshop"
    slide1.placeholders[1].text = "Building Insights from Complex Datasets\nQ2 2025 Training Series"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Workshop Agenda"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Session 1: Data Collection & Cleaning Strategies"
    for item in [
        "Session 2: Exploratory Data Analysis Techniques",
        "Session 3: Statistical Modeling Fundamentals",
        "Session 4: Visualization Best Practices",
        "Session 5: Real-World Case Studies",
        "Lunch Break: 12:00 PM - 1:00 PM",
        "Session 6: Hands-on Lab Exercise",
        "Session 7: Q&A and Wrap-up",
    ]:
        p = tf2.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 3: Introduction ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Why Data Analytics Matters"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Organizations using data-driven decisions are 23x more likely to acquire customers"
    for item in [
        "78% of leading firms consider analytics a competitive advantage",
        "Data literacy is now a core competency across all departments",
        "The analytics market is projected to reach $655B by 2029",
    ]:
        p = tf3.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 4: Data Collection ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Data Collection Strategies"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Primary Sources"
    items4 = [
        ("Surveys and questionnaires (SurveyMonkey, Typeform)", 1),
        ("API integrations (REST, GraphQL endpoints)", 1),
        ("IoT sensor data streams", 1),
        ("Secondary Sources", 0),
        ("Government databases (Census, BLS, WHO)", 1),
        ("Industry reports (Gartner, McKinsey, Deloitte)", 1),
        ("Open data repositories (Kaggle, UCI ML Repository)", 1),
    ]
    for text, level in items4:
        p = tf4.add_paragraph()
        p.text = text
        p.level = level

    # --- Slide 5: Data Cleaning ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Data Cleaning Pipeline"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Step 1: Identify missing values and anomalies"
    for item in [
        "Step 2: Handle duplicates and inconsistent formats",
        "Step 3: Normalize and standardize numeric fields",
        "Step 4: Validate referential integrity across tables",
        "Step 5: Document all transformations in a cleaning log",
        "Pro tip: Always preserve the raw dataset before cleaning",
    ]:
        p = tf5.add_paragraph()
        p.text = item

    # --- Slide 6: EDA Techniques ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Exploratory Data Analysis"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Descriptive Statistics"
    items6 = [
        ("Mean, median, mode, standard deviation", 1),
        ("Percentiles and interquartile range (IQR)", 1),
        ("Visualization Methods", 0),
        ("Histograms for distribution analysis", 1),
        ("Scatter plots for correlation detection", 1),
        ("Box plots for outlier identification", 1),
        ("Heatmaps for multivariate relationships", 1),
    ]
    for text, level in items6:
        p = tf6.add_paragraph()
        p.text = text
        p.level = level

    # --- Slide 7: Statistical Modeling ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Statistical Modeling Overview"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "Linear Regression: Predicting continuous outcomes"
    for item in [
        "Logistic Regression: Binary classification problems",
        "Decision Trees: Non-linear pattern recognition",
        "Random Forest: Ensemble learning for robust predictions",
        "Time Series: ARIMA, Prophet for temporal forecasting",
        "Clustering: K-means, DBSCAN for customer segmentation",
    ]:
        p = tf7.add_paragraph()
        p.text = item

    # --- Slide 8: Visualization Best Practices ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Visualization Best Practices"
    tf8 = slide8.placeholders[1].text_frame
    tf8.text = "Choose the right chart type for your data story"
    for item in [
        "Bar charts: Comparing categories (revenue by region)",
        "Line charts: Trends over time (monthly active users)",
        "Pie charts: Part-to-whole (market share, use sparingly)",
        "Treemaps: Hierarchical data (budget allocation)",
        "Always label axes, provide legends, use accessible colors",
        "Limit to 5-7 data series per visualization",
    ]:
        p = tf8.add_paragraph()
        p.text = item

    # --- Slide 9: Case Study ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Case Study: Meridian Retail Analytics"
    tf9 = slide9.placeholders[1].text_frame
    tf9.text = "Challenge: 15% decline in same-store sales over 18 months"
    for item in [
        "Approach: Integrated POS, CRM, and foot traffic data",
        "Key Finding: Peak shopping shifted from Saturday to Thursday",
        "Action: Reallocated staffing and promotions to Thursday-Friday",
        "Result: 22% revenue increase within one quarter",
        "ROI: $2.3M additional revenue on $180K analytics investment",
    ]:
        p = tf9.add_paragraph()
        p.text = item

    # --- Slide 10: Lab Exercise ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "Hands-on Lab: Customer Churn Analysis"
    tf10 = slide10.placeholders[1].text_frame
    tf10.text = "Dataset: Telecom customer records (7,043 rows, 21 features)"
    for item in [
        "Task 1: Load and profile the dataset using pandas",
        "Task 2: Identify top 5 predictors of churn",
        "Task 3: Build a logistic regression model",
        "Task 4: Evaluate with confusion matrix and ROC curve",
        "Task 5: Present findings in a 3-slide mini-deck",
        "Time allocation: 90 minutes",
    ]:
        p = tf10.add_paragraph()
        p.text = item

    # --- Slide 11: Wrap-up & Resources ---
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    slide11.shapes.title.text = "Resources & Next Steps"
    tf11 = slide11.placeholders[1].text_frame
    tf11.text = "Recommended Reading"
    items11 = [
        ("'Storytelling with Data' by Cole Nussbaumer Knaflic", 1),
        ("'Python for Data Analysis' by Wes McKinney", 1),
        ("'The Art of Statistics' by David Spiegelhalter", 1),
        ("Online Platforms", 0),
        ("Coursera: Google Data Analytics Certificate", 1),
        ("DataCamp: Interactive Python & R courses", 1),
        ("Contact: analytics-training@meridiangroup.com", 0),
    ]
    for text, level in items11:
        p = tf11.add_paragraph()
        p.text = text
        p.level = level

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
