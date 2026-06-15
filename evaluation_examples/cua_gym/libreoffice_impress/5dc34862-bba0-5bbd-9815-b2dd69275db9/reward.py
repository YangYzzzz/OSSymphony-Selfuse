"""
Reward Script: Create competitive comparison table on slide 7
Task ID: impress_sales_030
Domain: libreoffice_impress
Scoring:
  Component 1: Table exists on slide 7 with 5 rows x 4 columns (0.2 pts)
  Component 2: Header row correct (0.15 pts)
  Component 3: Data rows have correct cell values (0.3 pts)
  Component 4: 'Us' column cells (rows 1-4) colored green #00AA00 (0.35 pts)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_030'

# Expected table data
EXPECTED_HEADERS = ['Feature', 'Us', 'Competitor A', 'Competitor B']
EXPECTED_DATA = [
    ['Price', '$49/mo', '$79/mo', '$99/mo'],
    ['Users', 'Unlimited', '50', '25'],
    ['API Access', 'Yes', 'Limited', 'No'],
    ['Support', '24/7', 'Business hours', 'Email only'],
]


def persist_app_state():
    """Try to save any unsaved LibreOffice state."""
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 7 slides
    if len(prs.slides) < 7:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 7")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[6]  # slide 7, 0-indexed

    # Find the table shape on slide 7
    table_shape = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_shape = shape
            break

    # Component 1: Table exists on slide 7 with correct dimensions (0.2 points)
    try:
        if table_shape is None:
            print("FAIL: Component 1 -- No table found on slide 7")
        else:
            table = table_shape.table
            nrows = len(table.rows)
            ncols = len(table.columns)
            if nrows == 5 and ncols == 4:
                print(f"PASS: Component 1 -- Table found with {nrows}x{ncols} dimensions (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 1 -- Table dimensions are {nrows}x{ncols}, expected 5x4")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # If no table, remaining checks cannot proceed
    if table_shape is None:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    table = table_shape.table

    # Component 2: Header row has correct values (0.15 points)
    try:
        headers_match = True
        header_details = []
        for c in range(min(4, len(table.columns))):
            actual = table.cell(0, c).text.strip()
            expected = EXPECTED_HEADERS[c]
            if actual.lower() != expected.lower():
                headers_match = False
                header_details.append(f"col {c}: expected '{expected}', got '{actual}'")
            else:
                header_details.append(f"col {c}: '{actual}' OK")

        if headers_match:
            print(f"PASS: Component 2 -- Headers correct: {EXPECTED_HEADERS} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 -- Header mismatches: {header_details}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Data rows have correct values (0.3 points)
    # 4 data rows, each worth 0.075 points
    try:
        data_score = 0.0
        for r_idx, expected_row in enumerate(EXPECTED_DATA):
            row_num = r_idx + 1  # skip header row
            if row_num >= len(table.rows):
                print(f"FAIL: Component 3 -- Row {row_num} missing (table too short)")
                continue

            row_ok = True
            for c_idx, expected_val in enumerate(expected_row):
                if c_idx >= len(table.columns):
                    row_ok = False
                    continue
                actual = table.cell(row_num, c_idx).text.strip()
                if actual != expected_val:
                    row_ok = False
                    print(f"  FAIL: Cell ({row_num},{c_idx}): expected '{expected_val}', got '{actual}'")

            if row_ok:
                data_score += 0.075
                print(f"  PASS: Row {row_num} ({expected_row[0]}) data correct")
            else:
                print(f"  FAIL: Row {row_num} has mismatches")

        if data_score > 0:
            print(f"PASS: Component 3 -- Data rows partial score: {data_score:.3f}/0.3 pts")
            total_score += data_score
        else:
            print("FAIL: Component 3 -- No data rows correct")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Green (#00AA00) text in 'Us' column, rows 1-4 (0.35 points)
    # Each green cell worth 0.0875 points
    try:
        green_score = 0.0
        us_col = 1  # 'Us' is column index 1

        for r_idx in range(1, min(5, len(table.rows))):
            cell = table.cell(r_idx, us_col)
            cell_text = cell.text.strip()
            found_green = False

            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.color.type is not None:
                            rgb_str = str(run.font.color.rgb).upper()
                            if rgb_str == '00AA00':
                                found_green = True
                                break
                    except Exception:
                        pass
                if found_green:
                    break

            if found_green:
                green_score += 0.0875
                print(f"  PASS: Cell ({r_idx},{us_col}) '{cell_text}' has green #00AA00 color")
            else:
                print(f"  FAIL: Cell ({r_idx},{us_col}) '{cell_text}' missing green #00AA00 color")

        if green_score > 0:
            print(f"PASS: Component 4 -- Green color partial score: {green_score:.4f}/0.35 pts")
            total_score += green_score
        else:
            print("FAIL: Component 4 -- No green colored cells found in 'Us' column")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = round(min(total_score, 1.0), 4)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
