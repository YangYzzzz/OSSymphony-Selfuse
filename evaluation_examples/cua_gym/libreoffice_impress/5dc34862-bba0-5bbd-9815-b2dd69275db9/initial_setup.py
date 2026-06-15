"""
Initial Setup: Create competitive deck with 10 slides, slide 7 has title only
Task ID: impress_sales_030
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_030'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_body_slide(prs, title_text, body_lines):
    """Add a slide with layout 1 (Title + Content), set title and bullet body."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1].text_frame
    body.paragraphs[0].text = body_lines[0]
    for line in body_lines[1:]:
        p = body.add_paragraph()
        p.text = line
    return slide


def add_title_only_slide(prs, title_text):
    """Add a slide with layout 5 (Blank) and a title textbox."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)
    return slide


def create_initial():
    prs = Presentation()

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Competitive Deck 2025"
    slide1.placeholders[1].text = "Prepared by Strategic Sales Team | Q2 Review"

    # Slide 2: Executive Summary
    add_title_body_slide(prs, "Executive Summary", [
        "Our SaaS platform has achieved 34% YoY growth in enterprise accounts",
        "Net Promoter Score increased from 62 to 78 in the past quarter",
        "Three new product integrations launched since January",
        "Expansion into APAC markets underway with 12 signed partners",
        "Customer retention rate stands at 94.7%",
    ])

    # Slide 3: Our Product
    add_title_body_slide(prs, "Our Product", [
        "Cloud-native collaboration platform built for modern teams",
        "Real-time document editing with version history",
        "Advanced analytics dashboard with customizable widgets",
        "Seamless integration with 200+ third-party tools",
        "Enterprise-grade security with SOC 2 Type II certification",
        "Available on web, iOS, Android, and desktop",
    ])

    # Slide 4: Market Analysis
    add_title_body_slide(prs, "Market Analysis", [
        "Total addressable market: $18.3B by 2026 (Gartner)",
        "Collaboration tools segment growing at 14.2% CAGR",
        "Remote/hybrid work driving enterprise adoption",
        "Key verticals: Technology, Financial Services, Healthcare",
        "Mid-market segment ($10M-$500M revenue) most under-served",
        "Regulatory tailwinds in data sovereignty driving demand",
    ])

    # Slide 5: Key Features
    add_title_body_slide(prs, "Key Features", [
        "AI-powered workflow automation saves 12 hours per user monthly",
        "Unlimited project workspaces with granular permissions",
        "Built-in video conferencing with recording and transcription",
        "Custom API endpoints for enterprise integration needs",
        "Automated compliance reporting for regulated industries",
        "White-label options for partner channel distribution",
    ])

    # Slide 6: Pricing Strategy
    add_title_body_slide(prs, "Pricing Strategy", [
        "Starter: $29/mo — Individual professionals",
        "Professional: $49/mo — Small teams up to 20 users",
        "Enterprise: Custom — Unlimited users, dedicated support",
        "Annual billing discount: 20% across all tiers",
        "Free trial: 30 days with full feature access",
        "Volume licensing available for 500+ seat deployments",
    ])

    # Slide 7: How We Compare — TITLE ONLY, no table
    add_title_only_slide(prs, "How We Compare")

    # Slide 8: Customer Testimonials
    add_title_body_slide(prs, "Customer Testimonials", [
        '"Switching to this platform reduced our onboarding time by 60%"',
        "— Sarah Mitchell, VP of Operations at Meridian Health",
        "",
        '"The API flexibility lets us build exactly what our clients need"',
        "— David Park, CTO at FinBridge Solutions",
        "",
        '"Best support team in the industry, hands down"',
        "— Lisa Rodriguez, IT Director at Coastal Manufacturing",
    ])

    # Slide 9: Product Roadmap
    add_title_body_slide(prs, "Product Roadmap", [
        "Q2 2025: Advanced AI assistant with context-aware suggestions",
        "Q3 2025: Federated search across all connected platforms",
        "Q4 2025: Custom workflow builder with drag-and-drop interface",
        "Q1 2026: On-premise deployment option for regulated verticals",
        "Q2 2026: Multi-language real-time translation in collaboration",
    ])

    # Slide 10: Contact Us
    add_title_body_slide(prs, "Contact Us", [
        "Website: www.oursaasplatform.com",
        "Sales: sales@oursaasplatform.com | +1 (415) 555-0198",
        "Support: support@oursaasplatform.com | +1 (415) 555-0199",
        "Headquarters: 350 Mission Street, San Francisco, CA 94105",
        "Follow us: @OurSaaSPlatform on LinkedIn and Twitter",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
