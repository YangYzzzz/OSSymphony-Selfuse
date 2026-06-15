"""
FINAL REWARD SCRIPT - SUCCESS
Task: Insert a non-breaking space after 'Mr.', 'Ms.' and 'Dr.' in paragraph 3.
Generated: 2025-10-17 09:20:14
Status: success
Model: azure-o3
Total Steps: 9
"""

import os
import re
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

"""
Reward script for the task:
"Insert a non-breaking space after 'Mr.', 'Ms.' and 'Dr.' in paragraph 3."

The script verifies that, in paragraph 3 of the presentation, each of the three
specified abbreviations is immediately followed by a non-breaking space (U+00A0)
and **not** by a regular space.  A progressive score (0‒1) is awarded: 1⁄3 point
for each abbreviation that fulfils the requirement.  A perfect file therefore
receives 1.0.
"""

# ---------------------------------------------------------------------------
# CONFIGURATION
# ---------------------------------------------------------------------------
FILE_PATH = "/home/user/insert_a_non_breaking_space_after_mr_ms_and_dr_in_paragraph_3.pptx"
ABBREVIATIONS = ["Mr.", "Ms.", "Dr."]
NBSP = "\u00A0"  # non-breaking space character

# ---------------------------------------------------------------------------
# HELPER FUNCTIONS
# ---------------------------------------------------------------------------

def extract_paragraphs(prs):
    """Return a list of paragraph texts in document order (slide → shape → para)."""
    paragraphs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                para_text = "".join(run.text for run in para.runs)
                paragraphs.append(para_text)
    return paragraphs


def get_paragraph3(prs):
    """Locate paragraph 3 (1-based) from a non-title content placeholder."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            # Skip title / centre-title placeholders if present
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    continue
            if len(shape.text_frame.paragraphs) >= 3:
                p3 = shape.text_frame.paragraphs[2]
                return "".join(run.text for run in p3.runs)
    return None

# ---------------------------------------------------------------------------
# MAIN VERIFICATION LOGIC
# ---------------------------------------------------------------------------

def verify_non_breaking_space(file_path):
    total_score = 0.0
    per_abbr_score = 1.0 / len(ABBREVIATIONS)

    # 1. File existence
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    # 2. Load presentation
    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded presentation with {len(prs.slides)} slide(s)")
    except Exception as e:
        print(f"✗ Failed to load PPTX: {e}")
        print("REWARD: 0.0")
        return 0.0

    # 3. Retrieve paragraph 3
    para3 = get_paragraph3(prs)
    if para3 is None:
        print("✗ Could not locate paragraph 3 for verification")
        print("REWARD: 0.0")
        return 0.0

    print("Paragraph 3 text:", repr(para3))

    # 4. Check each abbreviation
    for abbr in ABBREVIATIONS:
        has_correct = (abbr + NBSP) in para3
        has_incorrect = (abbr + " ") in para3
        if has_correct and not has_incorrect:
            total_score += per_abbr_score
            print(f"✓ '{abbr}' uses non-breaking space (+{per_abbr_score:.2f})")
        else:
            if has_incorrect:
                print(f"✗ Regular space found after '{abbr}' – should be non-breaking")
            else:
                print(f"✗ '{abbr}' with required non-breaking space not found")

    final_score = round(min(total_score, 1.0), 2)
    print(f"Final score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score


# ---------------------------------------------------------------------------
# EXECUTION ENTRY-POINT
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    verify_non_breaking_space(FILE_PATH)
