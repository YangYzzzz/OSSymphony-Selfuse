"""
Initial Setup: Physics Equations teaching presentation with 8 slides.
Slide 5 has a title, explanatory text, and a text box with 'E = mc²' (no animations).
Task ID: impress_teach_047
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_047'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def add_equation_slide(prs, title_text, explanation_text, equation_text):
    """Add a slide with title, explanatory text, and an equation text box."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Title at top
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    # Explanatory text in middle-left area
    explain_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2.5))
    tf2 = explain_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = explanation_text
    run2 = p2.runs[0]
    run2.font.size = Pt(18)
    run2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Equation text box - prominent, centered
    eq_box = slide.shapes.add_textbox(Inches(2.5), Inches(4.2), Inches(5), Inches(1.5))
    tf3 = eq_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = equation_text
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.runs[0]
    run3.font.size = Pt(48)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0xC0, 0x39, 0x2B)

    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    add_title_slide(prs,
                    "Fundamental Physics Equations",
                    "Dr. Elena Vasquez — PHYS 301 Fall 2025")

    # Slide 2: Newton's Second Law
    add_equation_slide(prs,
                       "Newton's Second Law of Motion",
                       "The net force acting on an object equals the product of its mass and "
                       "acceleration. This foundational principle governs everything from "
                       "projectile motion to orbital mechanics.",
                       "F = ma")

    # Slide 3: Kinematic Equation
    add_equation_slide(prs,
                       "Kinematic Displacement Equation",
                       "This equation relates displacement to initial velocity, time, and "
                       "constant acceleration. It is essential for solving free-fall and "
                       "projectile problems in classical mechanics.",
                       "x = v₀t + ½at²")

    # Slide 4: Gravitational Force
    add_equation_slide(prs,
                       "Newton's Law of Universal Gravitation",
                       "Every particle in the universe attracts every other particle with a "
                       "force proportional to the product of their masses and inversely "
                       "proportional to the square of the distance between them.",
                       "F = Gm₁m₂ / r²")

    # Slide 5: Mass-Energy Equivalence (THE KEY SLIDE)
    add_equation_slide(prs,
                       "Mass-Energy Equivalence",
                       "Einstein's most famous equation demonstrates that mass and energy are "
                       "interchangeable. A small amount of mass can be converted into an "
                       "enormous amount of energy, as demonstrated by nuclear reactions.",
                       "E = mc²")

    # Slide 6: Schrodinger Equation
    add_equation_slide(prs,
                       "Time-Independent Schrödinger Equation",
                       "The cornerstone of quantum mechanics, this equation describes how "
                       "the quantum state of a physical system changes in space. Solutions "
                       "give the allowed energy levels of quantum systems.",
                       "Ĥψ = Eψ")

    # Slide 7: Maxwell's Equation (Gauss's Law)
    add_equation_slide(prs,
                       "Gauss's Law for Electricity",
                       "One of Maxwell's four equations, Gauss's law relates the electric "
                       "flux through a closed surface to the enclosed electric charge. It is "
                       "fundamental to electrostatics and electromagnetic theory.",
                       "∇ · E = ρ / ε₀")

    # Slide 8: Summary / Course Overview
    add_content_slide(prs,
                      "Summary & Next Steps",
                      ["These seven equations form the backbone of classical and modern physics.",
                       "Problem Set 3 due: September 29, 2025",
                       "Next lecture: Conservation of Energy and Momentum",
                       "Office hours: Tuesday & Thursday, 2:00–4:00 PM, Room 312B",
                       "Textbook reference: Halliday Ch. 7–9, Griffiths Ch. 1–2"])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the presentation
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
