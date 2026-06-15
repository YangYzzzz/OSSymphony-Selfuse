"""
Initial Setup: Create a 15-slide standard 4:3 presentation
Task ID: impress_fix_059
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_059'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                alignment=PP_ALIGN.LEFT, color=None, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_title(slide, text):
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.8),
                text, font_size=28, bold=True, color=RGBColor(0x1B, 0x4F, 0x72))


def add_bullets(slide, bullets, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(16)
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    bl = prs.slide_layouts[5]  # Blank

    # Slide 1: Title
    s = prs.slides.add_slide(bl)
    add_textbox(s, Inches(1.5), Inches(2.0), Inches(7), Inches(1.2),
                "Q4 2025 Business Review", font_size=36, bold=True,
                alignment=PP_ALIGN.CENTER, color=RGBColor(0x1B, 0x4F, 0x72))
    add_textbox(s, Inches(1.5), Inches(3.5), Inches(7), Inches(1.5),
                "Meridian Technologies Inc.\nPresented by Elena Vasquez, CFO",
                font_size=18, alignment=PP_ALIGN.CENTER, color=RGBColor(0x44, 0x44, 0x44))

    # Slide 2: Agenda
    s = prs.slides.add_slide(bl)
    add_title(s, "Agenda")
    add_bullets(s, [
        "1. Financial Performance Overview",
        "2. Revenue Breakdown by Region",
        "3. Product Line Analysis",
        "4. Customer Acquisition Metrics",
        "5. Operational Efficiency Updates",
        "6. Strategic Initiatives for 2026",
        "7. Q&A and Discussion",
    ], Inches(1), Inches(1.8), Inches(8), Inches(4.5))

    # Slide 3: Financial Highlights
    s = prs.slides.add_slide(bl)
    add_title(s, "Financial Highlights")
    for i, (value, label, change) in enumerate([
        ("$142.8M", "Total Revenue", "+18.3% YoY"),
        ("$31.4M", "Net Income", "+22.1% YoY"),
        ("67.2%", "Gross Margin", "+2.4pp"),
    ]):
        x = Inches(0.8 + i * 3.0)
        add_textbox(s, x, Inches(2.0), Inches(2.5), Inches(0.6),
                    value, font_size=28, bold=True, alignment=PP_ALIGN.CENTER,
                    color=RGBColor(0x1B, 0x4F, 0x72))
        add_textbox(s, x, Inches(2.6), Inches(2.5), Inches(0.4),
                    label, font_size=14, alignment=PP_ALIGN.CENTER,
                    color=RGBColor(0x33, 0x33, 0x33))
        add_textbox(s, x, Inches(3.0), Inches(2.5), Inches(0.4),
                    change, font_size=12, alignment=PP_ALIGN.CENTER,
                    color=RGBColor(0x2E, 0x7D, 0x32))

    # Slide 4: Revenue by Region
    s = prs.slides.add_slide(bl)
    add_title(s, "Revenue by Region")
    tbl = s.shapes.add_table(6, 4, Inches(1), Inches(2), Inches(8), Inches(3.5)).table
    for c, h in enumerate(["Region", "Q4 Revenue", "Q3 Revenue", "Change"]):
        tbl.cell(0, c).text = h
        for run in tbl.cell(0, c).text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row in enumerate([
        ["North America", "$62.4M", "$54.1M", "+15.3%"],
        ["Europe", "$38.7M", "$33.2M", "+16.6%"],
        ["Asia Pacific", "$27.1M", "$21.8M", "+24.3%"],
        ["Latin America", "$9.8M", "$8.4M", "+16.7%"],
        ["Middle East & Africa", "$4.8M", "$3.2M", "+50.0%"],
    ], 1):
        for c, val in enumerate(row):
            tbl.cell(r, c).text = val

    # Slide 5: Product Line Performance
    s = prs.slides.add_slide(bl)
    add_title(s, "Product Line Performance")
    for i, (name, rev, share, desc) in enumerate([
        ("CloudSync Pro", "$48.2M", "34%", "Enterprise cloud storage and sync"),
        ("DataVault Enterprise", "$35.6M", "25%", "Data warehousing solution"),
        ("SecureNet Gateway", "$28.9M", "20%", "Network security platform"),
        ("AnalytiQ Suite", "$18.3M", "13%", "Business intelligence tools"),
        ("DevOps Pipeline", "$11.8M", "8%", "CI/CD automation platform"),
    ]):
        y = Inches(1.8 + i * 1.0)
        add_textbox(s, Inches(0.8), y, Inches(3), Inches(0.4),
                    name, font_size=16, bold=True, color=RGBColor(0x1B, 0x4F, 0x72))
        add_textbox(s, Inches(4.0), y, Inches(1.5), Inches(0.4),
                    rev, font_size=14, alignment=PP_ALIGN.RIGHT)
        add_textbox(s, Inches(5.8), y, Inches(1.0), Inches(0.4),
                    share, font_size=14, alignment=PP_ALIGN.CENTER)
        add_textbox(s, Inches(7.0), y, Inches(2.5), Inches(0.4),
                    desc, font_size=11, color=RGBColor(0x66, 0x66, 0x66))

    # Slide 6: Customer Acquisition
    s = prs.slides.add_slide(bl)
    add_title(s, "Customer Acquisition Metrics")
    add_textbox(s, Inches(0.8), Inches(2.0), Inches(4), Inches(1.0),
                "New Enterprise Clients: 47", font_size=20, bold=True)
    add_textbox(s, Inches(0.8), Inches(3.0), Inches(4), Inches(1.0),
                "SMB Additions: 312", font_size=20, bold=True)
    add_textbox(s, Inches(5.5), Inches(2.0), Inches(4), Inches(1.0),
                "Customer Retention Rate: 94.7%", font_size=20, bold=True,
                color=RGBColor(0x2E, 0x7D, 0x32))
    add_textbox(s, Inches(5.5), Inches(3.0), Inches(4), Inches(1.0),
                "Net Promoter Score: 72", font_size=20, bold=True,
                color=RGBColor(0x1B, 0x4F, 0x72))
    add_textbox(s, Inches(0.8), Inches(4.5), Inches(8.5), Inches(1.5),
                "Enterprise pipeline grew 28% quarter-over-quarter, with particularly strong "
                "demand in the financial services and healthcare verticals. Average deal size "
                "increased to $285K from $241K in Q3.",
                font_size=13, color=RGBColor(0x55, 0x55, 0x55))

    # Slide 7: Operational Efficiency
    s = prs.slides.add_slide(bl)
    add_title(s, "Operational Efficiency")
    add_bullets(s, [
        "Cloud infrastructure costs reduced by 12% through optimization",
        "Average ticket resolution time improved to 4.2 hours (from 6.8)",
        "Automated deployment frequency increased to 47 per week",
        "System uptime maintained at 99.97% across all products",
        "Employee productivity index rose 8.5% after hybrid work policy",
    ], Inches(0.8), Inches(1.8), Inches(8.5), Inches(4.5))

    # Slide 8: Engineering Team Update
    s = prs.slides.add_slide(bl)
    add_title(s, "Engineering Team Update")
    tbl = s.shapes.add_table(5, 3, Inches(1.2), Inches(2.2), Inches(7.5), Inches(3)).table
    for c, h in enumerate(["Initiative", "Status", "ETA"]):
        tbl.cell(0, c).text = h
        for run in tbl.cell(0, c).text_frame.paragraphs[0].runs:
            run.font.bold = True
    for r, row in enumerate([
        ["Kubernetes Migration", "In Progress (78%)", "Q1 2026"],
        ["AI/ML Feature Integration", "Planning Phase", "Q2 2026"],
        ["API v3 Rollout", "Beta Testing", "Jan 2026"],
        ["Mobile App Redesign", "Completed", "Shipped Q4"],
    ], 1):
        for c, val in enumerate(row):
            tbl.cell(r, c).text = val

    # Slide 9: Marketing Campaign Results
    s = prs.slides.add_slide(bl)
    add_title(s, "Marketing Campaign Results")
    add_textbox(s, Inches(0.8), Inches(2.0), Inches(8.5), Inches(0.5),
                "Digital Marketing Performance - Q4 2025", font_size=16, bold=True,
                color=RGBColor(0x1B, 0x4F, 0x72))
    for i, (name, reach, engagement, cost) in enumerate([
        ("LinkedIn Thought Leadership", "2.4M impressions", "3.2% CTR", "$18.50 CPL"),
        ("Google Ads Enterprise", "890K impressions", "4.1% CTR", "$42.30 CPL"),
        ("Industry Conference Sponsorships", "15 events", "347 leads", "$125 CPL"),
        ("Content Syndication", "1.8M reach", "2.8% CTR", "$22.10 CPL"),
    ]):
        y = Inches(2.8 + i * 0.9)
        add_textbox(s, Inches(0.8), y, Inches(3.2), Inches(0.4), name, font_size=13, bold=True)
        add_textbox(s, Inches(4.2), y, Inches(1.8), Inches(0.4), reach, font_size=12)
        add_textbox(s, Inches(6.2), y, Inches(1.2), Inches(0.4), engagement, font_size=12)
        add_textbox(s, Inches(7.6), y, Inches(1.8), Inches(0.4), cost, font_size=12)

    # Slide 10: Partnership Ecosystem
    s = prs.slides.add_slide(bl)
    add_title(s, "Partnership Ecosystem")
    add_bullets(s, [
        "AWS Advanced Technology Partner - renewed with expanded scope",
        "Microsoft Azure Co-Sell agreement signed in October",
        "Salesforce AppExchange integration launched (4.6/5 rating)",
        "New SI partnerships: Deloitte Digital, Accenture Cloud",
        "Channel partner revenue grew 31% to $22.4M",
    ], Inches(0.8), Inches(1.8), Inches(8.5), Inches(4.5))

    # Slide 11: Risk Assessment
    s = prs.slides.add_slide(bl)
    add_title(s, "Risk Assessment")
    tbl = s.shapes.add_table(5, 3, Inches(1), Inches(2), Inches(8), Inches(3.5)).table
    for c, h in enumerate(["Risk Factor", "Severity", "Mitigation"]):
        tbl.cell(0, c).text = h
        for run in tbl.cell(0, c).text_frame.paragraphs[0].runs:
            run.font.bold = True
    for r, row in enumerate([
        ["Talent retention in AI/ML", "High", "Competitive comp packages, equity refresh"],
        ["Cloud cost escalation", "Medium", "FinOps team, reserved instance strategy"],
        ["Regulatory compliance (EU AI Act)", "Medium", "Legal review, compliance roadmap"],
        ["Cybersecurity threats", "High", "Zero-trust architecture, SOC expansion"],
    ], 1):
        for c, val in enumerate(row):
            tbl.cell(r, c).text = val

    # Slide 12: 2026 Strategic Priorities
    s = prs.slides.add_slide(bl)
    add_title(s, "2026 Strategic Priorities")
    for i, (title, desc) in enumerate([
        ("AI-First Product Strategy", "Embed machine learning across all product lines "
         "to deliver predictive analytics and automated workflows"),
        ("International Expansion", "Establish direct presence in Japan, Germany, and Brazil "
         "with localized go-to-market teams"),
        ("Platform Consolidation", "Unify CloudSync, DataVault, and AnalytiQ into a single "
         "integrated platform experience"),
        ("Sustainability Goals", "Achieve carbon-neutral operations by Q3 2026 and "
         "publish annual ESG report"),
    ]):
        y = Inches(1.8 + i * 1.3)
        add_textbox(s, Inches(0.8), y, Inches(8.5), Inches(0.4),
                    title, font_size=18, bold=True, color=RGBColor(0x1B, 0x4F, 0x72))
        add_textbox(s, Inches(0.8), Inches(1.8 + i * 1.3 + 0.45), Inches(8.5), Inches(0.6),
                    desc, font_size=13, color=RGBColor(0x44, 0x44, 0x44))

    # Slide 13: Financial Outlook
    s = prs.slides.add_slide(bl)
    add_title(s, "2026 Financial Outlook")
    for i, (metric, target, note) in enumerate([
        ("Revenue Target", "$168M - $175M", "+18-23% growth"),
        ("Operating Margin", "24% - 26%", "Up from 22.1%"),
        ("R&D Investment", "$38M - $42M", "26-28% of revenue"),
        ("Headcount Plan", "850 - 920 FTEs", "Current: 742"),
    ]):
        y = Inches(2.0 + i * 1.1)
        add_textbox(s, Inches(0.8), y, Inches(3.0), Inches(0.4),
                    metric, font_size=16, bold=True)
        add_textbox(s, Inches(4.2), y, Inches(2.5), Inches(0.4),
                    target, font_size=18, bold=True, color=RGBColor(0x1B, 0x4F, 0x72))
        add_textbox(s, Inches(7.0), y, Inches(2.5), Inches(0.4),
                    note, font_size=12, color=RGBColor(0x66, 0x66, 0x66))

    # Slide 14: Key Takeaways
    s = prs.slides.add_slide(bl)
    add_title(s, "Key Takeaways")
    add_bullets(s, [
        "Record Q4 revenue of $142.8M with strong margin expansion",
        "Customer base diversification reducing concentration risk",
        "Engineering velocity accelerating with modern infrastructure",
        "Clear strategic roadmap positions us for sustained growth",
        "Board-approved 2026 budget supports aggressive expansion",
    ], Inches(0.8), Inches(1.8), Inches(8.5), Inches(4.5))

    # Slide 15: Thank You
    s = prs.slides.add_slide(bl)
    add_textbox(s, Inches(2), Inches(2.5), Inches(6), Inches(1.0),
                "Thank You", font_size=40, bold=True, alignment=PP_ALIGN.CENTER,
                color=RGBColor(0x1B, 0x4F, 0x72))
    add_textbox(s, Inches(2), Inches(3.8), Inches(6), Inches(0.6),
                "Questions & Discussion", font_size=24, alignment=PP_ALIGN.CENTER,
                color=RGBColor(0x44, 0x44, 0x44))
    add_textbox(s, Inches(2), Inches(5.0), Inches(6), Inches(0.8),
                "Elena Vasquez, CFO\nevasquez@meridiantech.com\n+1 (415) 555-0142",
                font_size=14, alignment=PP_ALIGN.CENTER,
                color=RGBColor(0x66, 0x66, 0x66))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Slide dimensions: {prs.slide_width / 914400:.3f} x {prs.slide_height / 914400:.3f} inches')
    print(f'Number of slides: {len(prs.slides)}')

    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
