"""
Reward Script: Change slide size from standard (10x7.5) to widescreen (13.333x7.5)
Task ID: impress_fix_059
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Slide width matches widescreen 13.333 inches
  Component 2 (0.2): Slide count preserved at 15
  Component 3 (0.4): Content shapes have been scaled/repositioned for wider format
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_059'

# Constants
STANDARD_WIDTH_EMU = 9144000      # 10 inches
WIDESCREEN_WIDTH_EMU = 12192000   # 13.333 inches (approx)
EXPECTED_HEIGHT_EMU = 6858000     # 7.5 inches
EXPECTED_SLIDE_COUNT = 15
WIDTH_TOLERANCE = 0.005  # 0.5% relative tolerance


def persist_app_state(domain):
    """Save any unsaved GUI edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print(f"PERSIST: ctrl+s sent for {domain}")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def is_approx(val, expected, tol=WIDTH_TOLERANCE):
    if expected == 0:
        return val == 0
    return abs(val - expected) / abs(expected) <= tol


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide dimensions match widescreen 13.333 x 7.5 inches (0.4 points)
    try:
        actual_width = prs.slide_width
        actual_height = prs.slide_height
        width_ok = is_approx(actual_width, WIDESCREEN_WIDTH_EMU)
        height_ok = is_approx(actual_height, EXPECTED_HEIGHT_EMU)

        if width_ok and height_ok:
            print(f"PASS: Component 1 - Slide dimensions are widescreen "
                  f"({actual_width / 914400:.3f} x {actual_height / 914400:.3f} inches) (0.4 pts)")
            total_score += 0.4
        elif width_ok and not height_ok:
            print(f"PARTIAL: Component 1 - Width correct but height wrong "
                  f"({actual_width / 914400:.3f} x {actual_height / 914400:.3f} inches) (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 1 - Expected ~13.333 x 7.5 inches, "
                  f"found {actual_width / 914400:.3f} x {actual_height / 914400:.3f} inches")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Slide count preserved at 15 (0.2 points)
    # This is scored only in conjunction with the width change to avoid
    # awarding points on the initial file. The initial file also has 15 slides,
    # so we gate this on the width being widescreen.
    try:
        num_slides = len(prs.slides)
        width_is_widescreen = is_approx(prs.slide_width, WIDESCREEN_WIDTH_EMU)
        if num_slides == EXPECTED_SLIDE_COUNT and width_is_widescreen:
            print(f"PASS: Component 2 - Slide count is {num_slides} with widescreen width (0.2 pts)")
            total_score += 0.2
        elif num_slides == EXPECTED_SLIDE_COUNT and not width_is_widescreen:
            print(f"FAIL: Component 2 - Slide count is {num_slides} but width is not widescreen")
        else:
            print(f"FAIL: Component 2 - Expected {EXPECTED_SLIDE_COUNT} slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Content shapes have been scaled/repositioned (0.4 points)
    # On the initial file, shapes are laid out for 10-inch width.
    # On the golden file, non-placeholder TEXT_BOX shapes should have wider
    # positions/widths (scaled by ~1.333 factor). We check that non-placeholder
    # shapes on several slides have positions/widths that exceed the original
    # standard-width bounds, indicating scaling has occurred.
    # Specifically: any TEXT_BOX with left + width > STANDARD_WIDTH_EMU indicates
    # content extends beyond the original 10-inch boundary.
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        # Count shapes that extend beyond original 10-inch boundary on each slide
        slides_with_extended_content = 0
        slides_checked = 0

        for slide in prs.slides:
            textboxes = [s for s in slide.shapes
                         if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX]
            if not textboxes:
                continue
            slides_checked += 1
            # Check if any textbox has right edge beyond 9144000 (10 inches)
            has_extended = False
            for s in textboxes:
                right_edge = s.left + s.width
                if right_edge > STANDARD_WIDTH_EMU + 100000:  # at least ~0.1 inch beyond
                    has_extended = True
                    break
            if has_extended:
                slides_with_extended_content += 1

        if slides_checked == 0:
            print("FAIL: Component 3 - No text boxes found on any slide")
        else:
            ratio = slides_with_extended_content / slides_checked
            if ratio >= 0.7:
                print(f"PASS: Component 3 - {slides_with_extended_content}/{slides_checked} "
                      f"slides have content scaled beyond standard width (0.4 pts)")
                total_score += 0.4
            elif ratio >= 0.3:
                partial = round(0.4 * ratio, 2)
                print(f"PARTIAL: Component 3 - {slides_with_extended_content}/{slides_checked} "
                      f"slides scaled ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 3 - Only {slides_with_extended_content}/{slides_checked} "
                      f"slides have content beyond standard width")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
