"""
Reward Script: Insert pie chart with market share data on slide 2
Task ID: impress_rp_012
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Slide 2 contains a PIE chart
  Component 2 (0.25): Chart data values and categories match spec
  Component 3 (0.25): Chart slice colors match specified hex values
  Component 4 (0.25): Data labels show percentages and legend exists
"""

import os

from pptx import Presentation
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_012'

# Expected values from task spec
EXPECTED_CATEGORIES = ['Our Company', 'Competitor A', 'Competitor B', 'Others']
EXPECTED_VALUES = [35.0, 25.0, 20.0, 20.0]
EXPECTED_COLORS = {
    0: '2E86C1',  # Our Company - blue
    1: 'E67E22',  # Competitor A - orange
    2: '27AE60',  # Competitor B - green
    3: '95A5A6',  # Others - gray
}


def persist_app_state(domain):
    """Save any unsaved LibreOffice edits before verifying."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 2")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]

    # Find chart shape on slide 2
    chart_shape = None
    for shape in slide2.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Slide 2 contains a PIE chart (0.25 points)
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            # PIE chart type enum value is 5
            chart_type_val = chart.chart_type
            if str(chart_type_val) == 'PIE (5)' or 'PIE' in str(chart_type_val):
                print(f"PASS: Component 1 — Slide 2 has a PIE chart (type: {chart_type_val}) (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 — Chart found but type is {chart_type_val}, expected PIE")
        else:
            print("FAIL: Component 1 — No chart found on slide 2")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # If no chart, remaining components cannot pass
    if chart_shape is None:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    chart = chart_shape.chart
    plot = chart.plots[0]
    series = plot.series[0]

    # Component 2: Chart data values and categories match (0.25 points)
    try:
        actual_values = list(series.values)
        try:
            actual_categories = [str(c) for c in plot.categories]
        except Exception:
            actual_categories = []

        values_match = (actual_values == EXPECTED_VALUES)
        cats_match = (actual_categories == EXPECTED_CATEGORIES)

        if values_match and cats_match:
            print(f"PASS: Component 2 — Data values {actual_values} and categories {actual_categories} match (0.25 pts)")
            total_score += 0.25
        else:
            if not values_match:
                print(f"FAIL: Component 2 — Values mismatch: expected {EXPECTED_VALUES}, got {actual_values}")
            if not cats_match:
                print(f"FAIL: Component 2 — Categories mismatch: expected {EXPECTED_CATEGORIES}, got {actual_categories}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Chart slice colors match specified hex values (0.25 points)
    try:
        ns = {
            'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        }
        ser_xml = series._element
        dpts = ser_xml.findall('.//c:dPt', ns)

        colors_correct = 0
        colors_found = {}
        for dpt in dpts:
            idx_el = dpt.find('c:idx', ns)
            if idx_el is None:
                continue
            idx = int(idx_el.get('val'))
            srgb = dpt.find('.//a:srgbClr', ns)
            if srgb is not None:
                colors_found[idx] = srgb.get('val').upper()

        for idx, expected_color in EXPECTED_COLORS.items():
            actual_color = colors_found.get(idx, 'NOT_FOUND')
            if actual_color.upper() == expected_color.upper():
                colors_correct += 1
            else:
                print(f"  Color mismatch at index {idx}: expected {expected_color}, got {actual_color}")

        if colors_correct == 4:
            print(f"PASS: Component 3 — All 4 slice colors match spec (0.25 pts)")
            total_score += 0.25
        elif colors_correct >= 2:
            partial = round(0.25 * colors_correct / 4, 3)
            print(f"PARTIAL: Component 3 — {colors_correct}/4 colors match ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — Only {colors_correct}/4 colors match. Found: {colors_found}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Data labels show percentages + legend exists (0.25 points)
    try:
        sub_score = 0.0

        # Check data labels show percentage (via XML since API may report False)
        chart_xml = chart._chartSpace
        ns_c = {'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'}
        show_pct_els = chart_xml.findall('.//c:dLbls/c:showPercent', ns_c)
        has_pct_labels = any(el.get('val') == '1' for el in show_pct_els)

        if has_pct_labels:
            print(f"  Data labels show percentages: YES")
            sub_score += 0.15
        else:
            print(f"  Data labels show percentages: NO")

        # Check legend exists
        if chart.has_legend:
            print(f"  Legend present: YES")
            sub_score += 0.10
        else:
            print(f"  Legend present: NO")

        if sub_score > 0:
            print(f"PASS: Component 4 — Data labels/legend ({sub_score} pts)")
            total_score += sub_score
        else:
            print(f"FAIL: Component 4 — No data labels or legend found")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
