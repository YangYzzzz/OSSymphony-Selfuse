"""
Initial Setup: Create a 5-slide Market Analysis presentation with no charts.
Task ID: impress_rp_012
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_012'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q4 2025 Market Analysis Report"
    slide1.placeholders[1].text = "Prepared by Strategic Planning Division\nNovember 2025"

    # --- Slide 2: Market Share Overview (title + empty content) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content layout
    slide2.shapes.title.text = "Market Share Overview"
    # Leave the content placeholder empty — the agent must insert the pie chart here
    content_ph = slide2.placeholders[1]
    content_ph.text = ""  # explicitly empty

    # --- Slide 3: Revenue Breakdown ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Revenue Breakdown by Region"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = ""
    regions = [
        ("North America", "$12.4M", "+8.2% YoY"),
        ("Europe", "$8.7M", "+5.1% YoY"),
        ("Asia-Pacific", "$6.3M", "+14.7% YoY"),
        ("Latin America", "$2.1M", "+3.9% YoY"),
        ("Middle East & Africa", "$1.5M", "+6.8% YoY"),
    ]
    for i, (region, revenue, growth) in enumerate(regions):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = f"{region}: {revenue} ({growth})"
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(18)

    # --- Slide 4: Competitive Landscape ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Competitive Landscape"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = ""
    points = [
        "Our Company maintains market leadership with 35% share",
        "Competitor A gained 3 percentage points through aggressive pricing",
        "Competitor B focusing on premium segment with stable 20% share",
        "Fragmented 'Others' category represents emerging disruptors",
        "Key differentiator: superior customer retention (92% vs industry avg 78%)",
    ]
    for i, point in enumerate(points):
        if i == 0:
            p = tf4.paragraphs[0]
        else:
            p = tf4.add_paragraph()
        p.text = point
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(16)

    # --- Slide 5: Strategic Recommendations ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Strategic Recommendations"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = ""
    recs = [
        "Invest in AI-driven analytics to strengthen competitive moat",
        "Expand Asia-Pacific operations leveraging 14.7% growth trajectory",
        "Launch loyalty program to further improve retention metrics",
        "Explore strategic acquisition of emerging competitors in 'Others' segment",
    ]
    for i, rec in enumerate(recs):
        if i == 0:
            p = tf5.paragraphs[0]
        else:
            p = tf5.add_paragraph()
        p.text = rec
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(16)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
