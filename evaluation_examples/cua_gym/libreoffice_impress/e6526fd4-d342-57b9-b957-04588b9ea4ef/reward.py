"""
FINAL REWARD SCRIPT - SUCCESS
Task: I'm touching up slide 85 of my workshop deck and need to add a navigation control—a rounded rectangle button that measures exactly 6 cm by 1.8 cm and displays the word "Next". How do I set that up in LibreOffice Impress?
Generated: 2025-09-11 00:09:34
Status: success
Model: azure-o3
Total Steps: 1
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
import os, math

# -------------------------------------------------------------
# Reward Script: Verify "Next" Navigation Button on Slide 85
# -------------------------------------------------------------
# Task requirements:
# 1. On slide 85, there must be a rounded-rectangle shape.
# 2. The rounded rectangle must display the text "Next" (case-insensitive).
# 3. The shape must measure exactly 6 cm × 1.8 cm (±0.15 cm tolerance).
#
# Progressive scoring (total 1.0):
# • Correct text ............................................. 0.25
# • Shape is a rounded rectangle ............................. 0.25
# • Width within tolerance ................................... 0.25
# • Height within tolerance .................................. 0.25
#
# The script searches all shapes (including grouped shapes) on
# slide 85, evaluates any that contain the text "Next", and
# awards points according to the criteria above. The highest
# scoring candidate determines the final reward.
# -------------------------------------------------------------

EMU_PER_CM = 360000        # Conversion constant (English Metric Units → cm)
TOL_CM       = 0.15        # Dimensional tolerance in cm
TARGET_W_CM  = 6.0         # Required width  in cm
TARGET_H_CM  = 1.8         # Required height in cm
SLIDE_NUMBER = 85          # 1-based slide index


def _all_shapes(slide):
    """Return a flat list of all shapes (including children of groups)."""
    out = []

    def recurse(shp):
        out.append(shp)
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in shp.shapes:
                recurse(sub)
    for shp in slide.shapes:
        recurse(shp)
    return out


def verify_navigation_button(file_path,
                             slide_number=SLIDE_NUMBER,
                             target_w_cm=TARGET_W_CM,
                             target_h_cm=TARGET_H_CM,
                             tol_cm=TOL_CM):
    """Return a progressive score [0.0-1.0] verifying the navigation button."""
    max_score = 1.0
    best_candidate = 0.0

    # ------- Load presentation file -------
    if not os.path.exists(file_path):
        print("✗ File not found.")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Could not open presentation: {e}")
        return 0.0

    idx = slide_number - 1
    if idx < 0 or idx >= len(prs.slides):
        print(f"✗ Slide {slide_number} does not exist. Total slides: {len(prs.slides)}")
        return 0.0

    slide = prs.slides[idx]
    shapes = _all_shapes(slide)
    print(f"Total shapes found (incl. groups): {len(shapes)}")

    # ------- Evaluate candidate shapes -------
    for i, shp in enumerate(shapes):
        # Extract text (works for both placeholders & autoshapes)
        text = ""
        if hasattr(shp, 'text'):
            text = shp.text.strip()
        elif getattr(shp, 'has_text_frame', False):
            text = shp.text_frame.text.strip()

        if text.lower() != 'next':
            continue  # Not the navigation button

        print(f"\nEvaluating candidate shape #{i + 1} with text 'Next'")
        candidate = 0.25  # Correct text earns initial 0.25

        # Check for rounded rectangle
        is_round_rect = (
            shp.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and
            shp.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
        ) or ('roundRect' in shp.element.xml)

        if is_round_rect:
            candidate += 0.25
            print("  ✓ Rounded rectangle (+0.25)")
        else:
            print("  ✗ Not a rounded rectangle (+0)")

        # Dimensions check
        w_cm = shp.width  / EMU_PER_CM
        h_cm = shp.height / EMU_PER_CM
        print(f"  Shape size: {w_cm:.2f} cm × {h_cm:.2f} cm")

        if math.isclose(w_cm, target_w_cm, abs_tol=tol_cm):
            candidate += 0.25
            print("  ✓ Width within tolerance (+0.25)")
        else:
            print("  ✗ Width out of tolerance (+0)")

        if math.isclose(h_cm, target_h_cm, abs_tol=tol_cm):
            candidate += 0.25
            print("  ✓ Height within tolerance (+0.25)")
        else:
            print("  ✗ Height out of tolerance (+0)")

        print(f"  Candidate score: {candidate}")
        best_candidate = max(best_candidate, candidate)

    # ------- Final score -------
    final_score = min(best_candidate, max_score)
    print(f"\nBest candidate score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score


# -------------------------------------------------------------
# Execute verification when script is run directly
# -------------------------------------------------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/im_touching_up_slide_85_of_my_workshop_deck_and_need_to_add_a_navigation_controla_rounded_rectangle__golden.pptx"
    verify_navigation_button(FILE_PATH)

