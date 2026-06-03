"""
Initial Setup: Create a 20-slide training module presentation with one default master layout.
Task ID: impress_gf2_014
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_014'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_placeholder(placeholder, text, font_size=Pt(18), bold=False,
                             color=None, alignment=None):
    """Helper to set text on a placeholder with formatting."""
    placeholder.text = text
    for para in placeholder.text_frame.paragraphs:
        if alignment:
            para.alignment = alignment
        for run in para.runs:
            run.font.size = font_size
            run.font.bold = bold
            if color:
                run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=Pt(14),
                bold=False, color=None, alignment=None):
    """Add a textbox with formatted text to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()

    # Standard widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Module Training content - 20 slides covering various training topics
    slide_content = [
        {
            "layout": 0,  # Title Slide
            "title": "Module Training Program 2025",
            "subtitle": "Advanced Professional Development Series\nHuman Resources Division"
        },
        {
            "layout": 1,  # Title + Content
            "title": "Training Overview",
            "content": "This comprehensive training program covers essential skills for team leads and senior contributors. The program spans 12 weeks with bi-weekly assessment checkpoints."
        },
        {
            "layout": 1,
            "title": "Program Objectives",
            "content": "Develop leadership communication skills\nMaster project management fundamentals\nBuild cross-functional collaboration techniques\nStrengthen data-driven decision making"
        },
        {
            "layout": 1,
            "title": "Module 1: Leadership Communication",
            "content": "Effective communication is the cornerstone of successful leadership. This module covers active listening, constructive feedback delivery, and stakeholder management across distributed teams."
        },
        {
            "layout": 1,
            "title": "Communication Frameworks",
            "content": "STAR Method for structured feedback\nSituation-Behavior-Impact model\nNonviolent Communication principles\nActive listening techniques and practices"
        },
        {
            "layout": 1,
            "title": "Module 2: Project Management",
            "content": "Understanding project lifecycle from initiation through closure. Covers Agile and Waterfall methodologies with emphasis on hybrid approaches used in enterprise environments."
        },
        {
            "layout": 1,
            "title": "Agile Methodology Deep Dive",
            "content": "Sprint planning and estimation techniques\nDaily standup best practices\nRetrospective facilitation methods\nBacklog grooming and prioritization frameworks"
        },
        {
            "layout": 1,
            "title": "Risk Management",
            "content": "Identifying and categorizing project risks\nProbability-impact assessment matrices\nMitigation strategy development\nContingency planning and escalation paths"
        },
        {
            "layout": 1,
            "title": "Module 3: Cross-Functional Collaboration",
            "content": "Breaking down organizational silos requires intentional effort. Learn techniques for building trust across departments and aligning diverse teams toward shared objectives."
        },
        {
            "layout": 1,
            "title": "Stakeholder Mapping",
            "content": "Power-interest grid analysis\nInfluence strategy development\nCommunication cadence planning\nConflict resolution across team boundaries"
        },
        {
            "layout": 1,
            "title": "Module 4: Data-Driven Decisions",
            "content": "Modern leaders must be comfortable with data analysis. This module introduces key metrics, dashboard interpretation, and how to translate data insights into actionable strategies."
        },
        {
            "layout": 1,
            "title": "Key Performance Indicators",
            "content": "Revenue growth rate: 15% quarterly target\nCustomer satisfaction score: 4.5/5.0 minimum\nEmployee engagement index: 78% benchmark\nProject delivery rate: 92% on-time target"
        },
        {
            "layout": 1,
            "title": "Dashboard Interpretation",
            "content": "Understanding trend lines and seasonality patterns\nIdentifying statistical outliers vs normal variation\nCorrelation analysis between business metrics\nForecasting models and confidence intervals"
        },
        {
            "layout": 1,
            "title": "Module 5: Team Development",
            "content": "Building high-performing teams through structured development plans, mentoring programs, and performance coaching. Covers Tuckman's stages and psychological safety frameworks."
        },
        {
            "layout": 1,
            "title": "Performance Coaching",
            "content": "GROW coaching model implementation\nSetting SMART objectives with team members\nQuarterly review conversation templates\nCareer development path planning"
        },
        {
            "layout": 1,
            "title": "Assessment Schedule",
            "content": "Week 2: Leadership Communication Assessment\nWeek 4: Project Management Simulation\nWeek 6: Collaboration Case Study\nWeek 8: Data Analysis Challenge\nWeek 10: Team Development Portfolio\nWeek 12: Capstone Presentation"
        },
        {
            "layout": 1,
            "title": "Resources and Materials",
            "content": "All training materials are available on the company learning platform. Supplementary reading lists and practice exercises are distributed at the start of each module."
        },
        {
            "layout": 1,
            "title": "Support Contacts",
            "content": "Program Director: Sarah Chen (sarah.chen@company.com)\nLead Facilitator: Marcus Rivera (marcus.rivera@company.com)\nTechnical Support: helpdesk@company.com\nHR Liaison: Priya Patel (priya.patel@company.com)"
        },
        {
            "layout": 1,
            "title": "Certification Requirements",
            "content": "Minimum 80% attendance across all modules\nCompletion of all six assessments with passing scores\nSubmission of capstone project by Week 12 deadline\nPeer evaluation participation for at least 3 colleagues"
        },
        {
            "layout": 1,
            "title": "Next Steps",
            "content": "Review the pre-work materials for Module 1\nComplete the self-assessment survey by Friday\nJoin the dedicated Slack channel #module-training-2025\nSchedule your first 1:1 with your assigned mentor"
        },
    ]

    for i, sc in enumerate(slide_content):
        layout_idx = sc["layout"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        if slide.shapes.title:
            slide.shapes.title.text = sc["title"]

        if layout_idx == 0 and "subtitle" in sc:
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = sc["subtitle"]
        elif "content" in sc:
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = sc["content"]

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
