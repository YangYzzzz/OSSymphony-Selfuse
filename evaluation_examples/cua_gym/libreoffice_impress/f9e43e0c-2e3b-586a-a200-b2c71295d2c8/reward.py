"""
Reward Script: Verify master slide layouts for Section Divider and Two Column Content
Task ID: impress_gf2_014
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): 'Section Divider' layout exists by name
  Component 2 (0.20): Section Divider has dark background (#1E293B)
  Component 3 (0.15): Section Divider has title placeholder only (no content placeholders)
  Component 4 (0.20): 'Two Column Content' layout exists by name
  Component 5 (0.10): Two Column Content has white background (#FFFFFF)
  Component 6 (0.15): Two Column Content has title + two content placeholders arranged side-by-side
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_014'


def persist_app_state(domain):
    """Try to save any unsaved LibreOffice edits."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print(f"PERSIST: ctrl+s sent for {domain}")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def find_layout_by_name(prs, name):
    """Find a slide layout by name across all masters."""
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    return None


def get_non_meta_placeholders(layout):
    """Return placeholders that are NOT date/footer/slide-number (idx < 10)."""
    return [ph for ph in layout.placeholders if ph.placeholder_format.idx < 10]


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # ---- Component 1: 'Section Divider' layout exists (0.20 points) ----
    sd_layout = None
    try:
        sd_layout = find_layout_by_name(prs, 'Section Divider')
        if sd_layout is not None:
            print(f"PASS: Component 1 -- 'Section Divider' layout found (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 -- No layout named 'Section Divider' found")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # ---- Component 2: Section Divider background is #1E293B (0.20 points) ----
    try:
        if sd_layout is not None:
            fill = sd_layout.background.fill
            if fill.type is not None and fill.type == 1:  # SOLID fill
                bg_color = str(fill.fore_color.rgb).upper()
                if bg_color == '1E293B':
                    print(f"PASS: Component 2 -- Section Divider background is #1E293B (0.20 pts)")
                    total_score += 0.20
                else:
                    print(f"FAIL: Component 2 -- Section Divider background is #{bg_color}, expected #1E293B")
            else:
                print(f"FAIL: Component 2 -- Section Divider background fill type is {fill.type}, expected SOLID (1)")
        else:
            print(f"FAIL: Component 2 -- Skipped (no Section Divider layout)")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # ---- Component 3: Section Divider has ONLY title placeholder, no content (0.15 points) ----
    try:
        if sd_layout is not None:
            non_meta_phs = get_non_meta_placeholders(sd_layout)
            # Should have exactly 1 placeholder (title) with no content/body placeholders
            has_title = False
            has_content = False
            for ph in non_meta_phs:
                ph_type = ph.placeholder_format.type
                # TITLE (1) or CENTER_TITLE (3) count as title
                if ph_type in (1, 3):
                    has_title = True
                else:
                    has_content = True

            if has_title and not has_content:
                print(f"PASS: Component 3 -- Section Divider has title only, no content placeholders (0.15 pts)")
                total_score += 0.15
            elif not has_title:
                print(f"FAIL: Component 3 -- Section Divider missing title placeholder")
            else:
                print(f"FAIL: Component 3 -- Section Divider has content placeholders (expected title only)")
        else:
            print(f"FAIL: Component 3 -- Skipped (no Section Divider layout)")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # ---- Component 4: 'Two Column Content' layout exists (0.20 points) ----
    tc_layout = None
    try:
        tc_layout = find_layout_by_name(prs, 'Two Column Content')
        if tc_layout is not None:
            print(f"PASS: Component 4 -- 'Two Column Content' layout found (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 4 -- No layout named 'Two Column Content' found")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # ---- Component 5: Two Column Content has white background (0.10 points) ----
    try:
        if tc_layout is not None:
            fill = tc_layout.background.fill
            if fill.type is not None and fill.type == 1:  # SOLID fill
                bg_color = str(fill.fore_color.rgb).upper()
                if bg_color == 'FFFFFF':
                    print(f"PASS: Component 5 -- Two Column Content background is #FFFFFF (0.10 pts)")
                    total_score += 0.10
                else:
                    print(f"FAIL: Component 5 -- Two Column Content background is #{bg_color}, expected #FFFFFF")
            else:
                # White could also be no fill / inherited - check if fill type is None (transparent/default)
                # Default slide background is typically white, so no explicit fill could be acceptable
                if fill.type is None:
                    print(f"PASS: Component 5 -- Two Column Content has no explicit fill (defaults to white) (0.10 pts)")
                    total_score += 0.10
                else:
                    print(f"FAIL: Component 5 -- Two Column Content fill type is {fill.type}, expected SOLID white or no fill")
        else:
            print(f"FAIL: Component 5 -- Skipped (no Two Column Content layout)")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # ---- Component 6: Two Column Content has title + two side-by-side content placeholders (0.15 points) ----
    try:
        if tc_layout is not None:
            non_meta_phs = get_non_meta_placeholders(tc_layout)
            has_title = False
            content_phs = []
            for ph in non_meta_phs:
                ph_type = ph.placeholder_format.type
                if ph_type in (1, 3):  # TITLE or CENTER_TITLE
                    has_title = True
                elif ph_type in (7, 2):  # OBJECT or BODY — content placeholder
                    content_phs.append(ph)

            if has_title and len(content_phs) == 2:
                # Check that the two content placeholders are side-by-side (similar top, different left)
                ph1, ph2 = content_phs
                tops_similar = abs(ph1.top - ph2.top) < Emu(914400)  # within 1 inch
                lefts_different = abs(ph1.left - ph2.left) > Emu(914400)  # more than 1 inch apart
                # Check roughly equal widths (within 20% of each other)
                if ph1.width > 0 and ph2.width > 0:
                    width_ratio = min(ph1.width, ph2.width) / max(ph1.width, ph2.width)
                    widths_similar = width_ratio > 0.7
                else:
                    widths_similar = False

                if tops_similar and lefts_different and widths_similar:
                    print(f"PASS: Component 6 -- Two Column Content has title + 2 equal side-by-side columns (0.15 pts)")
                    total_score += 0.15
                else:
                    print(f"FAIL: Component 6 -- Columns not properly side-by-side "
                          f"(tops_similar={tops_similar}, lefts_different={lefts_different}, widths_similar={widths_similar})")
            elif not has_title:
                print(f"FAIL: Component 6 -- Two Column Content missing title placeholder")
            else:
                print(f"FAIL: Component 6 -- Two Column Content has {len(content_phs)} content placeholders, expected 2")
        else:
            print(f"FAIL: Component 6 -- Skipped (no Two Column Content layout)")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
