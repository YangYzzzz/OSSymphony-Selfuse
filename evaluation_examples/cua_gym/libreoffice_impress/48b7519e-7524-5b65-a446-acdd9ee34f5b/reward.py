"""
Reward Script: Animated organizational chart on slide 4
Task ID: impress_stu_063
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): 7 org chart boxes with correct text labels on slide 4
  Component 2 (0.20): Correct fill colors per level (blue/teal/green)
  Component 3 (0.15): At least 6 connecting lines (connectors) on slide 4
  Component 4 (0.25): Fade animations with 3 click-triggered groups
  Component 5 (0.15): Box size hierarchy (President > VP > third-row)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_063'

# Expected org chart labels
LEVEL1_LABELS = ['President']
LEVEL2_LABELS = ['VP Academic', 'VP Social']
LEVEL3_LABELS = ['Academic Events', 'Tutoring', 'Events', 'Communications']
ALL_LABELS = LEVEL1_LABELS + LEVEL2_LABELS + LEVEL3_LABELS


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # slide 4 (0-indexed)

    # Collect shapes by type
    auto_shapes = []  # org chart boxes
    connectors = []   # connecting lines
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            text = shape.text.strip() if hasattr(shape, 'text') else ''
            auto_shapes.append((shape, text))
        elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
            connectors.append(shape)

    # =========================================================
    # Component 1: 7 org chart boxes with correct text (0.25 pts)
    # =========================================================
    try:
        found_labels = [text for (_, text) in auto_shapes if text]
        matched_labels = [label for label in ALL_LABELS if label in found_labels]
        match_ratio = len(matched_labels) / len(ALL_LABELS) if ALL_LABELS else 0

        if match_ratio >= 1.0:
            print(f"PASS: Component 1 - All 7 org chart labels found: {matched_labels} (0.25 pts)")
            total_score += 0.25
        elif match_ratio >= 0.5:
            partial = round(0.25 * match_ratio, 3)
            print(f"PARTIAL: Component 1 - {len(matched_labels)}/7 labels found: {matched_labels} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 - Only {len(matched_labels)}/7 labels found: {matched_labels}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # =========================================================
    # Component 2: Correct fill colors per level (0.20 pts)
    # =========================================================
    try:
        color_points = 0.0
        color_checks = 0
        color_total = 3  # 3 levels to check

        # Build a map of label -> fill color
        label_color_map = {}
        for shape, text in auto_shapes:
            if text in ALL_LABELS:
                try:
                    fill = shape.fill
                    if fill.type is not None and fill.type == 1:  # solid fill
                        label_color_map[text] = str(fill.fore_color.rgb).upper()
                except Exception:
                    pass

        # Check level 1: President should be blue-ish
        pres_color = label_color_map.get('President', '')
        if pres_color:
            # Accept blue hues - check R < G and R < B, B is dominant or close
            r, g, b = int(pres_color[0:2], 16), int(pres_color[2:4], 16), int(pres_color[4:6], 16)
            if b > r and b > g:
                color_checks += 1
                print(f"PASS: Level 1 color - President has blue fill ({pres_color})")
            else:
                print(f"FAIL: Level 1 color - President fill {pres_color} is not blue (expected blue-dominant)")
        else:
            print(f"FAIL: Level 1 color - President has no solid fill color")

        # Check level 2: VP boxes should be teal-ish
        vp_ok = 0
        for label in LEVEL2_LABELS:
            c = label_color_map.get(label, '')
            if c:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                if g >= r and b >= r:  # teal = green+blue dominant
                    vp_ok += 1
        if vp_ok == len(LEVEL2_LABELS):
            color_checks += 1
            print(f"PASS: Level 2 color - VP boxes have teal fill")
        else:
            print(f"FAIL: Level 2 color - {vp_ok}/{len(LEVEL2_LABELS)} VP boxes have teal fill")

        # Check level 3: Third-row should be green-ish
        l3_ok = 0
        for label in LEVEL3_LABELS:
            c = label_color_map.get(label, '')
            if c:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                if g > r:  # green dominant over red
                    l3_ok += 1
        if l3_ok == len(LEVEL3_LABELS):
            color_checks += 1
            print(f"PASS: Level 3 color - Third-row boxes have green fill")
        else:
            print(f"FAIL: Level 3 color - {l3_ok}/{len(LEVEL3_LABELS)} third-row boxes have green fill")

        color_score = round(0.20 * (color_checks / color_total), 3)
        if color_score > 0:
            print(f"PASS: Component 2 - {color_checks}/{color_total} color levels correct ({color_score} pts)")
            total_score += color_score
        else:
            print(f"FAIL: Component 2 - No correct color levels")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # =========================================================
    # Component 3: At least 6 connecting lines (0.15 pts)
    # =========================================================
    try:
        num_connectors = len(connectors)
        if num_connectors >= 6:
            print(f"PASS: Component 3 - {num_connectors} connectors found (>= 6 required) (0.15 pts)")
            total_score += 0.15
        elif num_connectors >= 3:
            partial = round(0.15 * (num_connectors / 6), 3)
            print(f"PARTIAL: Component 3 - {num_connectors}/6 connectors found ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 - Only {num_connectors} connectors found (need >= 6)")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # =========================================================
    # Component 4: Fade animations with 3 click groups (0.25 pts)
    # =========================================================
    try:
        ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        fade_count = 0
        click_groups = 0

        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide4.xml') as f:
                content = f.read().decode()
                root = ET.fromstring(content)

                # Find all animation par elements
                # presetID="10" = Fade, presetClass="entr" = entrance
                # nodeType="clickEffect" = new click group
                # nodeType="withEffect" = same click group
                all_ns = {
                    'p': ns_p,
                    'a': ns_a,
                }

                # Parse all cTn elements with presetID
                anim_elements = root.findall('.//{%s}cTn' % ns_p)
                for elem in anim_elements:
                    preset_id = elem.get('presetID')
                    preset_class = elem.get('presetClass')
                    node_type = elem.get('nodeType')
                    if preset_id and preset_class == 'entr':
                        fade_count += 1 if preset_id == '10' else 0
                        if node_type == 'clickEffect':
                            click_groups += 1

        if click_groups >= 3 and fade_count >= 3:
            print(f"PASS: Component 4 - {click_groups} click groups, {fade_count} fade animations (0.25 pts)")
            total_score += 0.25
        elif click_groups >= 2 and fade_count >= 2:
            partial = round(0.25 * 0.6, 3)
            print(f"PARTIAL: Component 4 - {click_groups} click groups (need 3), {fade_count} fade anims ({partial} pts)")
            total_score += partial
        elif click_groups >= 1 and fade_count >= 1:
            partial = round(0.25 * 0.3, 3)
            print(f"PARTIAL: Component 4 - {click_groups} click groups, {fade_count} fade anims ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 - No fade animations with click groups found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # =========================================================
    # Component 5: Box size hierarchy (0.15 pts)
    # President (largest) > VP (medium) > third-row (smallest)
    # =========================================================
    try:
        # Collect areas by level
        level_areas = {1: [], 2: [], 3: []}
        for shape, text in auto_shapes:
            if text in LEVEL1_LABELS:
                level_areas[1].append(shape.width * shape.height)
            elif text in LEVEL2_LABELS:
                level_areas[2].append(shape.width * shape.height)
            elif text in LEVEL3_LABELS:
                level_areas[3].append(shape.width * shape.height)

        hierarchy_ok = 0
        hierarchy_total = 2  # level1 > level2, and level2 > level3

        if level_areas[1] and level_areas[2]:
            avg1 = sum(level_areas[1]) / len(level_areas[1])
            avg2 = sum(level_areas[2]) / len(level_areas[2])
            if avg1 > avg2:
                hierarchy_ok += 1
                print(f"PASS: Level 1 > Level 2 area ({avg1:.0f} > {avg2:.0f})")
            else:
                print(f"FAIL: Level 1 area ({avg1:.0f}) not > Level 2 ({avg2:.0f})")

        if level_areas[2] and level_areas[3]:
            avg2 = sum(level_areas[2]) / len(level_areas[2])
            avg3 = sum(level_areas[3]) / len(level_areas[3])
            if avg2 > avg3:
                hierarchy_ok += 1
                print(f"PASS: Level 2 > Level 3 area ({avg2:.0f} > {avg3:.0f})")
            else:
                print(f"FAIL: Level 2 area ({avg2:.0f}) not > Level 3 ({avg3:.0f})")

        if hierarchy_ok > 0:
            size_score = round(0.15 * (hierarchy_ok / hierarchy_total), 3)
            print(f"PASS: Component 5 - {hierarchy_ok}/{hierarchy_total} size hierarchy checks ({size_score} pts)")
            total_score += size_score
        else:
            print(f"FAIL: Component 5 - Size hierarchy not satisfied")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
persist_app_state()
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
