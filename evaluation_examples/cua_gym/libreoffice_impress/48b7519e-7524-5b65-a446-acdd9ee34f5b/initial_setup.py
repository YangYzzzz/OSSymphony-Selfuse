"""
Initial Setup: Create a 7-slide Student Government presentation with empty slide 4.
Task ID: impress_stu_063
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_063'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a styled text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullet_slide(slide, title_text, bullets):
    """Add title and bullet points to a slide."""
    # Title
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                title_text, font_size=32, bold=True,
                color=RGBColor(0x1A, 0x3C, 0x6E), alignment=PP_ALIGN.LEFT)
    # Bullets
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # Background
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
    # Title
    add_textbox(slide1, Inches(1), Inches(2), Inches(8), Inches(1.5),
                "Student Government Association", font_size=36, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    # Subtitle
    add_textbox(slide1, Inches(1.5), Inches(3.8), Inches(7), Inches(1),
                "Annual Presentation 2025-2026", font_size=24, bold=False,
                color=RGBColor(0xCC, 0xDD, 0xFF), alignment=PP_ALIGN.CENTER)
    # School name
    add_textbox(slide1, Inches(2), Inches(5), Inches(6), Inches(0.6),
                "Westfield University", font_size=20, bold=False,
                color=RGBColor(0xAA, 0xBB, 0xDD), alignment=PP_ALIGN.CENTER)

    # --- Slide 2: Our Mission ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide2, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Our Mission", font_size=32, bold=True,
                color=RGBColor(0x1A, 0x3C, 0x6E), alignment=PP_ALIGN.LEFT)
    mission_text = (
        "The Student Government Association at Westfield University is dedicated to "
        "representing the interests of the entire student body. We serve as the primary "
        "liaison between students and university administration, advocating for policies "
        "that enhance academic excellence, campus life, and student welfare.\n\n"
        "Our commitment extends to fostering an inclusive community where every student's "
        "voice is heard. Through transparent governance, collaborative leadership, and "
        "proactive engagement, we strive to create meaningful change that benefits all "
        "members of our campus community."
    )
    add_textbox(slide2, Inches(0.8), Inches(1.5), Inches(8.4), Inches(5),
                mission_text, font_size=18, bold=False,
                color=RGBColor(0x33, 0x33, 0x33), alignment=PP_ALIGN.LEFT)

    # --- Slide 3: Key Initiatives ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide3, "Key Initiatives", [
        "• Mental Health Awareness Week — partnering with counseling services for campus-wide workshops",
        "• Sustainability Drive — introducing recycling stations in all dormitories by Fall 2025",
        "• Academic Support Expansion — launching free peer tutoring for STEM courses",
        "• Cultural Festival — celebrating diversity with food, music, and art from 30+ countries",
        "• Transportation Improvement — negotiating extended shuttle hours for evening classes",
        "• Digital Accessibility — upgrading campus Wi-Fi infrastructure in study areas",
    ])

    # --- Slide 4: Organization Structure (EMPTY - task target) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide4, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Organization Structure", font_size=32, bold=True,
                color=RGBColor(0x1A, 0x3C, 0x6E), alignment=PP_ALIGN.LEFT)
    # Body intentionally left empty — this is where the agent must create the org chart

    # --- Slide 5: Upcoming Events ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide5, "Upcoming Events", [
        "• Sep 12 — Welcome Back BBQ at the Student Center lawn (3:00 PM - 7:00 PM)",
        "• Sep 28 — Club Fair featuring 85+ registered student organizations",
        "• Oct 15 — Town Hall Meeting with Dean Martinez on curriculum changes",
        "• Nov 3 — Fall Formal Dance at the Grand Ballroom ($15 tickets on sale Oct 1)",
        "• Nov 20 — Thanksgiving Food Drive collection deadline",
        "• Dec 5 — End-of-Semester Study Break with free snacks and therapy dogs",
    ])

    # --- Slide 6: Budget Overview ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide6, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Budget Overview", font_size=32, bold=True,
                color=RGBColor(0x1A, 0x3C, 0x6E), alignment=PP_ALIGN.LEFT)
    # Table
    rows, cols = 7, 3
    table_shape = slide6.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(4))
    table = table_shape.table
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(2.25)
    table.columns[2].width = Inches(2.25)

    headers = ["Category", "Allocated", "Spent"]
    data_rows = [
        ["Student Events", "$28,500", "$19,340"],
        ["Academic Programs", "$15,200", "$11,750"],
        ["Marketing & Outreach", "$8,000", "$6,210"],
        ["Equipment & Supplies", "$5,500", "$3,890"],
        ["Emergency Fund", "$3,000", "$800"],
        ["Total", "$60,200", "$41,990"],
    ]

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    for r, row_data in enumerate(data_rows, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(14)
                if r == len(data_rows):  # Total row
                    run.font.bold = True

    # --- Slide 7: Contact Us ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide7, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Contact Us", font_size=32, bold=True,
                color=RGBColor(0x1A, 0x3C, 0x6E), alignment=PP_ALIGN.LEFT)
    contact_info = [
        "• Office: Student Center, Room 204",
        "• Hours: Monday–Friday, 9:00 AM – 5:00 PM",
        "• Email: sga@westfield.edu",
        "• Phone: (555) 234-8901",
        "• Instagram: @westfield_sga",
        "• Website: sga.westfield.edu",
    ]
    add_bullet_slide(slide7, "Contact Us", contact_info)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
