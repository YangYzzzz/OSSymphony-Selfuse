"""
FINAL REWARD SCRIPT - SUCCESS
Task: While polishing slide 21, I noticed it’s missing a footer number. In LibreOffice Impress, how do I drop the slide number smack in the bottom-center and tint it Gray 50% (#808080)?
Generated: 2025-09-10 22:00:08
Status: success
Model: azure-o3
Total Steps: 6
"""

import os
from pptx import Presentation

def verify_slide_footer_number(file_path: str,
                                target_slide_number: int = 21,
                                tolerance_ratio: float = 0.05) -> float:
    """Verify that slide `target_slide_number` contains a slide-number shape
    centred at the bottom of the slide and coloured Gray 50 % (#808080).

    Returns a progressive score from 0.0 to 1.0.
    """

    print(f"Verifying presentation file: {file_path}")
    max_score = 1.0
    total_score = 0.0

    # ------------------------------------------------------------------
    # 0.  Basic checks (NO points for mere existence!)
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File not found.")
        return 0.0
    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded presentation with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        return 0.0

    # ------------------------------------------------------------------
    # 1.  Target slide exists (0.2 pts)
    # ------------------------------------------------------------------
    if len(prs.slides) >= target_slide_number:
        total_score += 0.2
        print("✓ Presentation contains the target slide (0.2 points)")
    else:
        print("✗ Presentation lacks the target slide – stopping verification")
        return total_score  # cannot continue

    slide = prs.slides[target_slide_number - 1]
    slide_w, slide_h = prs.slide_width, prs.slide_height

    # ------------------------------------------------------------------
    # 2.  Find shapes whose text equals the slide number
    # ------------------------------------------------------------------
    candidates = []
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        text = shape.text.strip()
        if text == str(target_slide_number):
            center_x = shape.left + shape.width / 2
            top = shape.top
            rgb_val = None
            # Safely attempt to read RGB colour
            try:
                run = shape.text_frame.paragraphs[0].runs[0]
                rgb = run.font.color.rgb  # may be None
                if rgb is not None:
                    rgb_val = tuple(rgb)  # (R, G, B)
            except Exception:
                pass
            candidates.append({"center_x": center_x, "top": top, "rgb": rgb_val})
            print(f"Found candidate shape at (center_x={center_x}, top={top}) with RGB {rgb_val}")

    if not candidates:
        print("✗ No shape with the correct slide number text found")
        return total_score

    # Slide-number text present ⇒ +0.2
    total_score += 0.2
    print("✓ Slide-number text present on target slide (0.2 points)")

    # ------------------------------------------------------------------
    # 3.  Evaluate each candidate for position & colour
    # ------------------------------------------------------------------
    center_tolerance = slide_w * tolerance_ratio  # allowable deviation from centre
    bottom_threshold = slide_h * (1 - 0.20)      # bottom 20 % of slide

    position_awarded = False
    colour_awarded = False

    for cand in candidates:
        # a) Position check – centred horizontally & within bottom 20 %
        if not position_awarded:
            is_centered = abs(cand["center_x"] - slide_w / 2) <= center_tolerance
            is_bottom   = cand["top"] >= bottom_threshold
            if is_centered and is_bottom:
                position_awarded = True
                total_score += 0.3
                print("✓ Slide number positioned bottom-centre (0.3 points)")

        # b) Colour check – Gray 50 % (#808080)
        if not colour_awarded and cand["rgb"] == (128, 128, 128):
            colour_awarded = True
            total_score += 0.3
            print("✓ Slide-number colour is Gray 50 % (#808080) (0.3 points)")

        if position_awarded and colour_awarded:
            break  # no need to inspect other candidates

    # Diagnostic messages for missing aspects
    if not position_awarded:
        print("✗ No candidate meets bottom-centre positioning requirement")
    if not colour_awarded:
        print("✗ No candidate uses Gray 50 % colour (#808080)")

    final_score = min(max_score, total_score)
    print(f"Total score: {final_score} / {max_score}")
    return final_score

# ----------------------------------------------------------------------
# Execute verification when run as a script
# ----------------------------------------------------------------------
if __name__ == "__main__":
    PRESENTATION_PATH = (
        "/home/user/while_polishing_slide_21_i_noticed_its_missing_a_footer_number_"
        "in_libreoffice_impress_how_do_i_drop__golden.pptx"
    )

    reward = verify_slide_footer_number(PRESENTATION_PATH)
    print(f"REWARD: {reward}")
