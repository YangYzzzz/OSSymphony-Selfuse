"""
Reward Script: Save presentation as ODP while keeping original PPTX unchanged
Task ID: impress_gf3_046
Domain: libreoffice_impress
Scoring:
  Component 1: ODP file exists at /home/user/Documents/Presentation_ODP_Version.odp (0.15)
  Component 2: ODP is valid OpenDocument Presentation format (0.20)
  Component 3: ODP contains 8 slides (0.25)
  Component 4: ODP slide titles match original PPTX content (0.20)
  Component 5: Original PPTX unchanged AND ODP exists (compound check) (0.20)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_046'

ODP_PATH = '/home/user/Documents/Presentation_ODP_Version.odp'
PPTX_PATH = '/home/user/Original_Deck.pptx'

# Expected slide titles from the original presentation (ground truth)
EXPECTED_TITLES = [
    'Q4 2025 Strategic Review',
    'Agenda',
    'Revenue Performance',
    'Client Acquisition Pipeline',
    'Product Development Roadmap',
    'Regional Expansion Strategy',
    'Financial Outlook 2026',
    'Key Actions & Next Steps',
]


def get_odp_slide_info(odp_path):
    """Parse ODP file and return list of (slide_name, [texts]) tuples."""
    ns = {
        'draw': 'urn:oasis:names:tc:opendocument:xmlns:drawing:1.0',
        'text': 'urn:oasis:names:tc:opendocument:xmlns:text:1.0',
    }
    slides = []
    with zipfile.ZipFile(odp_path, 'r') as zf:
        with zf.open('content.xml') as f:
            root = ET.parse(f).getroot()
            pages = root.findall('.//draw:page', ns)
            for page in pages:
                name = page.get('{urn:oasis:names:tc:opendocument:xmlns:drawing:1.0}name', '')
                texts = []
                for tp in page.findall('.//text:p', ns):
                    t = ''.join(tp.itertext()).strip()
                    if t:
                        texts.append(t)
                slides.append((name, texts))
    return slides


def verify_task():
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0
    odp_found = os.path.exists(ODP_PATH) and os.path.isfile(ODP_PATH)

    # Component 1: ODP file exists at correct path (0.15 points)
    try:
        if odp_found and os.path.getsize(ODP_PATH) > 0:
            print(f"PASS: Component 1 -- ODP file exists at {ODP_PATH} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 -- ODP file not found or empty at {ODP_PATH}")
            odp_found = False
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")
        odp_found = False

    if not odp_found:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: ODP is valid OpenDocument Presentation (0.20 points)
    try:
        with zipfile.ZipFile(ODP_PATH, 'r') as zf:
            names = zf.namelist()
            has_content = 'content.xml' in names
            has_manifest = 'META-INF/manifest.xml' in names
            mimetype_correct = False
            if 'mimetype' in names:
                mt = zf.read('mimetype').decode().strip()
                mimetype_correct = (mt == 'application/vnd.oasis.opendocument.presentation')

            if has_content and has_manifest and mimetype_correct:
                print(f"PASS: Component 2 -- Valid ODP format (content.xml, manifest, mimetype correct) (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 2 -- Invalid ODP: content.xml={has_content}, manifest={has_manifest}, mimetype={mimetype_correct}")
    except zipfile.BadZipFile:
        print(f"FAIL: Component 2 -- File is not a valid ZIP/ODP archive")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: ODP contains exactly 8 slides (0.25 points)
    try:
        slides = get_odp_slide_info(ODP_PATH)
        slide_count = len(slides)
        if slide_count == 8:
            print(f"PASS: Component 3 -- ODP has 8 slides (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 -- Expected 8 slides, found {slide_count}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: ODP slide titles match original PPTX content (0.20 points)
    try:
        slides = get_odp_slide_info(ODP_PATH)
        matches = 0
        for i, expected_title in enumerate(EXPECTED_TITLES):
            if i < len(slides):
                slide_name, texts = slides[i]
                # Check if the expected title appears in the slide name or text content
                all_text = ' '.join([slide_name] + texts)
                if expected_title in all_text:
                    matches += 1
                else:
                    print(f"  INFO: Slide {i+1} missing title '{expected_title}', found: {slide_name}, texts: {texts[:2]}")

        if matches == 8:
            print(f"PASS: Component 4 -- All 8 slide titles match original content (0.20 pts)")
            total_score += 0.20
        elif matches >= 6:
            partial = round(0.20 * matches / 8, 2)
            print(f"PARTIAL: Component 4 -- {matches}/8 slide titles match ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 -- Only {matches}/8 slide titles match")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Original PPTX still valid with 8 slides AND ODP exists (0.20 points)
    # This is a compound check: only scored when ODP exists (ensures it differentiates initial vs golden)
    try:
        from pptx import Presentation
        prs = Presentation(PPTX_PATH)
        pptx_slide_count = len(prs.slides)
        # Check first slide title to confirm content unchanged
        first_title = ''
        for shape in prs.slides[0].shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = p.text.strip()
                    if t:
                        first_title = t
                        break
                if first_title:
                    break

        if pptx_slide_count == 8 and 'Q4 2025 Strategic Review' in first_title:
            print(f"PASS: Component 5 -- Original PPTX intact: {pptx_slide_count} slides, title='{first_title}' (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 5 -- PPTX modified: slides={pptx_slide_count}, title='{first_title}'")
    except Exception as e:
        print(f"ERROR: Component 5 -- Cannot load original PPTX: {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
if not os.path.exists(PPTX_PATH):
    print(f"CRITICAL: Original PPTX not found: {PPTX_PATH}")
    print("REWARD: 0.0")
else:
    verify_task()
