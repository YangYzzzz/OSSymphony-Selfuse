"""
Initial Setup: Create an 8-slide presentation and open in LibreOffice Impress
Task ID: impress_gf3_046
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_046'
OUTPUT = f'{WORKDIR}/Original_Deck.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q4 2025 Strategic Review"
    slide1.placeholders[1].text = "Meridian Analytics Group\nPrepared by the Strategy Division"

    # ---- Slide 2: Agenda ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Revenue Performance & Market Trends"
    for item in [
        "Client Acquisition Pipeline",
        "Product Development Roadmap",
        "Regional Expansion Strategy",
        "Financial Outlook & Projections",
        "Key Risks and Mitigation Plans",
        "Q&A and Next Steps",
    ]:
        p = tf2.add_paragraph()
        p.text = item
        p.level = 0

    # ---- Slide 3: Revenue Overview ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Revenue Performance"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Total Revenue: $48.7M (+12.3% YoY)"
    items3 = [
        "Enterprise Segment: $29.1M (59.8%)",
        "SMB Segment: $12.4M (25.5%)",
        "Government Contracts: $7.2M (14.7%)",
        "Recurring revenue ratio improved to 73% from 68%",
    ]
    for item in items3:
        p = tf3.add_paragraph()
        p.text = item
        p.level = 0

    # ---- Slide 4: Client Acquisition ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Client Acquisition Pipeline"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "New Clients Won: 34 (Target: 30)"
    items4 = [
        "Healthcare Vertical: 12 new accounts",
        "Financial Services: 9 new accounts",
        "Technology: 8 new accounts",
        "Manufacturing: 5 new accounts",
        "Average Deal Size: $142K (up from $118K)",
    ]
    for item in items4:
        p = tf4.add_paragraph()
        p.text = item
        p.level = 0

    # ---- Slide 5: Product Roadmap ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Product Development Roadmap"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Platform v3.2 Launch: January 2026"
    items5 = [
        "AI-powered analytics dashboard (Beta: Nov 2025)",
        "Real-time collaboration features",
        "Enhanced API gateway with OAuth 2.0 support",
        "Mobile app redesign (iOS + Android)",
        "SOC 2 Type II certification in progress",
    ]
    for item in items5:
        p = tf5.add_paragraph()
        p.text = item
        p.level = 0

    # ---- Slide 6: Regional Expansion ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Regional Expansion Strategy"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Current Presence: North America, Western Europe"
    items6 = [
        "APAC Office Opening: Singapore (Q1 2026)",
        "LATAM Partnership: Sao Paulo-based reseller agreement",
        "DACH Region: Dedicated sales team hired",
        "Target: 25% international revenue by 2027",
    ]
    for item in items6:
        p = tf6.add_paragraph()
        p.text = item
        p.level = 0

    # ---- Slide 7: Financial Outlook ----
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Financial Outlook 2026"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "Projected Revenue: $58.4M (+19.9%)"
    items7 = [
        "EBITDA Target: $14.6M (25% margin)",
        "Headcount Growth: 180 → 220 employees",
        "R&D Investment: $8.2M (14% of revenue)",
        "Customer Retention Target: 94%",
        "Cash Runway: 18+ months at current burn",
    ]
    for item in items7:
        p = tf7.add_paragraph()
        p.text = item
        p.level = 0

    # ---- Slide 8: Next Steps ----
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Key Actions & Next Steps"
    tf8 = slide8.placeholders[1].text_frame
    tf8.text = "Board Review: December 15, 2025"
    items8 = [
        "Finalize Singapore office lease by November 30",
        "Complete SOC 2 audit by end of Q4",
        "Launch AI dashboard beta to 50 pilot customers",
        "Submit Series C term sheet for review",
        "Schedule all-hands meeting for January 8, 2026",
    ]
    for item in items8:
        p = tf8.add_paragraph()
        p.text = item
        p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Ensure no ODP file exists
    odp_path = f'{WORKDIR}/Documents/Presentation_ODP_Version.odp'
    if os.path.exists(odp_path):
        os.remove(odp_path)

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
