"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 94 is giving me layout headaches. In LibreOffice Impress, how can I push the title all the way to the left edge and have the body text box sit on the right side, right-aligned but still ragged (no full justification)?
Generated: 2025-09-10 16:53:22
Status: success
Model: azure-o3
Total Steps: 7
"""

import os
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

"""
Reward Script for LibreOffice Impress Task
Task recap:
Slide 94 should have …
  1. Title textbox pushed completely to the left edge
  2. Body textbox pushed to the right edge (right-flush position)
  3. Body paragraphs right-aligned (ragged right, NOT fully-justified)
The script awards:
  • 0.4 pts  – Title flush left
  • 0.3 pts  – Body box flush right
  • 0.3 pts  – Body paragraphs right-aligned
Progressive scoring ensures partial credit.
A perfect slide gives REWARD = 1.0.
"""

def load_presentation(path: str):
    """Load a .pptx file and return Presentation object or None."""
    if not os.path.exists(path):
        print(f"✗ File not found: {path}")
        return None
    try:
        prs = Presentation(path)
        print(f"✓ Loaded presentation with {len(prs.slides)} slides")
        return prs
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        return None


def is_title_shape(shape):
    """Identify title / subtitle placeholders."""
    return (
        shape.is_placeholder
        and shape.placeholder_format.type
        in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.SUBTITLE}
    )


def find_body_shape(slide):
    """Find the body/content textbox on the slide."""
    # 1) common body-like placeholder types
    for shp in slide.shapes:
        if shp.is_placeholder and shp.placeholder_format.type in {
            PP_PLACEHOLDER.BODY,
            PP_PLACEHOLDER.OBJECT,
            PP_PLACEHOLDER.VERTICAL_BODY,
        }:
            return shp
    # 2) any text shape that isn't the title
    for shp in slide.shapes:
        if hasattr(shp, "text_frame") and not is_title_shape(shp):
            return shp
    return None


def verify_title_left(slide, tolerance=91440):
    """Check that title left-edge ≤ tolerance (≈0.1 inch)."""
    for shp in slide.shapes:
        if is_title_shape(shp):
            print(f"Title left = {shp.left} EMUs")
            if shp.left <= tolerance:
                print("✓ Title is flush to left edge")
                return True
            print("✗ Title is not flush left")
            return False
    print("✗ No title placeholder found")
    return False


def verify_body_position(slide, slide_width, tolerance=91440):
    """Check that body textbox right-edge ≈ slide width."""
    body = find_body_shape(slide)
    if not body:
        print("✗ Body textbox not found")
        return False
    right_edge = body.left + body.width
    print(f"Body right edge = {right_edge}  (Slide width = {slide_width})")
    if abs(right_edge - slide_width) <= tolerance:
        print("✓ Body textbox is flush with right edge")
        return True
    print("✗ Body textbox is not flush right")
    return False


def verify_body_alignment(slide):
    """Ensure every paragraph in body textbox is RIGHT-aligned (ragged)."""
    body = find_body_shape(slide)
    if not body or not hasattr(body, "text_frame"):
        print("✗ Body textbox not suitable for alignment check")
        return False
    tf = body.text_frame
    if not tf.paragraphs:
        print("✗ No paragraphs in body textbox")
        return False
    for idx, p in enumerate(tf.paragraphs, 1):
        if p.alignment != PP_PARAGRAPH_ALIGNMENT.RIGHT:
            print(f"✗ Paragraph {idx} alignment is not RIGHT (found {p.alignment})")
            return False
    print("✓ All body paragraphs are right-aligned (ragged)")
    return True


def verify_task(file_path: str) -> float:
    """Main verification routine – returns score ∈ [0, 1]."""
    score = 0.0
    prs = load_presentation(file_path)
    if not prs:
        return 0.0

    # Ensure slide 94 exists (index 93)
    if len(prs.slides) < 94:
        print(f"✗ Presentation only has {len(prs.slides)} slides – need ≥ 94")
        return 0.0
    slide = prs.slides[93]
    slide_width = prs.slide_width

    # 1) Title flush left (0.4 pts)
    if verify_title_left(slide):
        score += 0.4

    # 2) Body box flush right (0.3 pts)
    if verify_body_position(slide, slide_width):
        score += 0.3

    # 3) Body paragraphs right-aligned (0.3 pts)
    if verify_body_alignment(slide):
        score += 0.3

    final_score = round(min(score, 1.0), 4)
    print(f"Total score: {final_score}")
    return final_score


if __name__ == "__main__":
    FILE = (
        "/home/user/slide_94_is_giving_me_layout_headaches_in_libreoffice_impress_"
        "how_can_i_push_the_title_all_the_way_t_golden.pptx"
    )
    reward = verify_task(FILE)
    print(f"REWARD: {reward}")

