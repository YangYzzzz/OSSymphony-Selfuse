"""
FINAL REWARD SCRIPT - SUCCESS
Task: On Slide 1 in LibreOffice Impress I’ve got three separate text boxes (they’re currently all plain white). I want to give them a traffic-light look: set Text Box 1’s fill to #FFFF00 yellow, Text Box 2’s fill to #FF0000 red, and Text Box 3’s fill to #00FF00 green. What’s the quickest way to apply those exact colors?
Generated: 2025-09-10 13:03:02
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def verify_traffic_light_coloring(file_path: str) -> float:
    """Verify that slide 1 has three text boxes filled with the
    exact traffic-light colors: yellow (#FFFF00), red (#FF0000),
    and green (#00FF00).  Progressive scoring is used:
        1/3 per correctly colored text box (max 1.0).
    """

    expected_colors = ["FFFF00", "FF0000", "00FF00"]  # hex without '#'
    max_score = 1.0
    matched = 0

    # ---- Preliminary checks (no points given) ----
    if not os.path.exists(file_path):
        print("✗ Presentation file not found.")
        print("REWARD: 0.0")
        return 0.0
    if not file_path.lower().endswith(".pptx"):
        print("✗ Unsupported file type – expected .pptx file.")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as exc:
        print(f"✗ Could not open presentation: {exc}")
        print("REWARD: 0.0")
        return 0.0

    if not prs.slides:
        print("✗ Presentation contains no slides.")
        print("REWARD: 0.0")
        return 0.0

    print(f"✓ Presentation loaded – {len(prs.slides)} slide(s) found.")

    # ---- Inspect shapes on the first slide ----
    slide = prs.slides[0]
    found_colors = []

    for idx, shape in enumerate(slide.shapes):
        # Consider only text boxes (shape_type TEXT_BOX) or shapes with a text_frame
        is_text_box = (
            shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or hasattr(shape, "text_frame")
        )
        if not is_text_box:
            continue

        fill = shape.fill
        rgb_hex = None
        # 1 == MSO_FILL.SOLID; solid fill is needed to read fore_color.rgb
        if fill.type == 1 and fill.fore_color.rgb is not None:
            rgb_hex = str(fill.fore_color.rgb).upper()

        if rgb_hex:
            print(f"  Shape {idx}: Detected solid fill #{rgb_hex}")
            found_colors.append(rgb_hex)
        else:
            print(f"  Shape {idx}: No solid fill colour detected")

    # ---- Match detected colours against expected colours ----
    remaining = found_colors.copy()  # avoid mutating original list in debug prints
    for exp in expected_colors:
        if exp in remaining:
            matched += 1
            remaining.remove(exp)  # prevent double-counting same colour
            print(f"✓ Found expected colour #{exp}")
        else:
            print(f"✗ Missing expected colour #{exp}")

    # ---- Progressive scoring ----
    score = (matched / len(expected_colors)) * max_score
    score = round(min(max_score, score), 4)  # limit to 1.0 & tidy precision

    print(f"\nMatched colours: {matched}/{len(expected_colors)}")
    print(f"REWARD: {score}")
    return score


# ---------------------------
# Execute verification (only when run as script)
# ---------------------------
if __name__ == "__main__":
    FILE = "/home/user/on_slide_1_in_libreoffice_impress_ive_got_three_separate_text_boxes_theyre_currently_all_plain_white_golden.pptx"
    verify_traffic_light_coloring(FILE)
