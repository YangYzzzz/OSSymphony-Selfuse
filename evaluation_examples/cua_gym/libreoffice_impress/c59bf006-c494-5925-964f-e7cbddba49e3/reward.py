"""
Reward Script: Move summary table on slide 3 to the bottom and resize to full width
Task ID: osworld_impress_table_position_bottom_008
Domain: libreoffice_impress
Scoring:
  - Component 1: Table left position aligns with content area left margin (0.4 pts)
  - Component 2: Table width spans full content area width (0.4 pts)
  - Component 3: Table is repositioned to the bottom of the slide (0.2 pts)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_table_position_bottom_008'

# Slide dimensions (10 x 7.5 inches in EMU)
SLIDE_WIDTH_EMU = 9144000
SLIDE_HEIGHT_EMU = 6858000

# Content area bounds (matching other shapes on the slide)
CONTENT_LEFT_EMU = 457200    # left edge of content area
CONTENT_WIDTH_EMU = 8229600  # full content width

# Initial table position (center-right placement before task)
INITIAL_TABLE_LEFT = 5303520
INITIAL_TABLE_TOP = 1371600

# Tolerance: 0.5% relative tolerance for position/size checks
TOLERANCE = 0.005

# Threshold: table must be positioned at or below 60% of slide height to be "at bottom"
BOTTOM_THRESHOLD = int(SLIDE_HEIGHT_EMU * 0.50)


def is_approx_equal(val1, val2, tolerance=TOLERANCE):
    """Check if two values are approximately equal with relative tolerance."""
    if val1 == val2:
        return True
    if val1 == 0 or val2 == 0:
        return val1 == val2
    return abs(val1 - val2) / max(abs(val1), abs(val2)) <= tolerance


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"CRITICAL: Expected at least 3 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]  # 0-indexed, slide 3

    # Find the table shape on slide 3
    table_shape = None
    for shape in slide3.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_shape = shape
            break

    if table_shape is None:
        print("CRITICAL: No table found on slide 3")
        print("REWARD: 0.0")
        return 0.0

    table_left = table_shape.left
    table_top = table_shape.top
    table_width = table_shape.width
    table_height = table_shape.height

    print(f"Table on slide 3: left={table_left}, top={table_top}, width={table_width}, height={table_height}")
    print(f"Expected (golden): left={CONTENT_LEFT_EMU}, top>=bottom_threshold({BOTTOM_THRESHOLD}), width={CONTENT_WIDTH_EMU}")

    # Component 1: Table left position aligns with content area left margin (0.4 points)
    # FAILS on initial (left=5303520) → PASSES on golden (left=457200)
    try:
        if is_approx_equal(table_left, CONTENT_LEFT_EMU):
            print(f"PASS: Component 1 — Table left aligns with content area ({table_left} ≈ {CONTENT_LEFT_EMU}) (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 — Table left={table_left}, expected ~{CONTENT_LEFT_EMU} (content area left margin)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Table width spans full content area width (0.4 points)
    # FAILS on initial (width=3474720) → PASSES on golden (width=8229600)
    try:
        if is_approx_equal(table_width, CONTENT_WIDTH_EMU):
            print(f"PASS: Component 2 — Table spans full content width ({table_width} ≈ {CONTENT_WIDTH_EMU}) (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 2 — Table width={table_width}, expected ~{CONTENT_WIDTH_EMU} (full content width)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Table top is positioned at the bottom of the slide (0.2 points)
    # FAILS on initial (top=1371600) → PASSES on golden (top=4297680, > 50% of slide height)
    try:
        if table_top >= BOTTOM_THRESHOLD:
            print(f"PASS: Component 3 — Table top={table_top} is at bottom of slide (>= {BOTTOM_THRESHOLD} threshold) (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 3 — Table top={table_top}, expected >= {BOTTOM_THRESHOLD} (bottom half of slide)")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the VM environment
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
