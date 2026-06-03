"""
Initial Setup: 5-slide executive summary deck with table on slide 3 in center-right area
Task ID: osworld_impress_table_position_bottom_008
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_table_position_bottom_008'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_box(slide, text, left, top, width, height, font_size=28, bold=True, color=RGBColor(0x1F, 0x49, 0x7D)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.bold = bold
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    return txBox


def add_body_box(slide, text, left, top, width, height, font_size=16, color=RGBColor(0x33, 0x33, 0x33)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    # Standard widescreen 10x7.5 inches
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Background
    fill1 = slide1.background.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    add_title_box(slide1, "Meridian Capital Partners",
                  Inches(1), Inches(2.5), Inches(8), Inches(1.2),
                  font_size=36, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_body_box(slide1, "Executive Summary — Q1 2025 Investor Briefing",
                 Inches(1), Inches(3.8), Inches(8), Inches(0.8),
                 font_size=20, color=RGBColor(0xC8, 0xD8, 0xEB))
    add_body_box(slide1, "Confidential | March 2025",
                 Inches(1), Inches(6.5), Inches(4), Inches(0.5),
                 font_size=12, color=RGBColor(0xC8, 0xD8, 0xEB))

    # ---- Slide 2: Investment Overview ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)

    add_title_box(slide2, "Investment Overview",
                  Inches(0.5), Inches(0.4), Inches(9), Inches(0.7),
                  font_size=28, bold=True, color=RGBColor(0x1F, 0x49, 0x7D))

    # Horizontal divider (thin rectangle)
    rect = slide2.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.5), Inches(1.15), Inches(9), Emu(36000)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x1F, 0x49, 0x7D)
    rect.line.fill.background()

    content2 = (
        "Meridian Capital Partners manages $2.4B in assets across three primary verticals: "
        "growth equity, venture lending, and real assets. Our Q1 2025 portfolio demonstrated "
        "resilient performance with aggregate returns of 12.3% against a benchmark of 8.7%.\n\n"
        "Key highlights:\n"
        "  • Portfolio companies: 47 active investments\n"
        "  • New commitments this quarter: $180M\n"
        "  • Exits completed: 3 (total proceeds $340M)\n"
        "  • Pipeline under review: $620M"
    )
    add_body_box(slide2, content2,
                 Inches(0.5), Inches(1.4), Inches(9), Inches(5.5),
                 font_size=15, color=RGBColor(0x22, 0x22, 0x22))

    # ---- Slide 3: Financial Performance (with summary table in center-right) ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill3 = slide3.background.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)

    add_title_box(slide3, "Financial Performance — Q1 2025",
                  Inches(0.5), Inches(0.4), Inches(9), Inches(0.7),
                  font_size=28, bold=True, color=RGBColor(0x1F, 0x49, 0x7D))

    rect3 = slide3.shapes.add_shape(
        1,
        Inches(0.5), Inches(1.15), Inches(9), Emu(36000)
    )
    rect3.fill.solid()
    rect3.fill.fore_color.rgb = RGBColor(0x1F, 0x49, 0x7D)
    rect3.line.fill.background()

    # Left body text
    content3 = (
        "Revenue for Q1 2025 reached $84.2M, representing 18% year-over-year growth. "
        "EBITDA margins expanded by 210 basis points to 34.5%, driven by operational "
        "efficiencies in the growth equity segment.\n\n"
        "Operating expenses were tightly managed at $55.1M, with technology investments "
        "accounting for 22% of total OpEx. Free cash flow generation improved to $28.6M."
    )
    add_body_box(slide3, content3,
                 Inches(0.5), Inches(1.4), Inches(5.2), Inches(3.0),
                 font_size=14, color=RGBColor(0x22, 0x22, 0x22))

    # Summary table — positioned in center-right area (NOT at bottom, NOT full width)
    # Center-right: left=5.8", top=1.5", width=3.8", height=2.8"
    table_left = Inches(5.8)
    table_top = Inches(1.5)
    table_width = Inches(3.8)
    table_height = Inches(2.8)

    table_shape = slide3.shapes.add_table(
        5, 2,
        table_left, table_top,
        table_width, table_height
    )
    tbl = table_shape.table

    # Column widths
    tbl.columns[0].width = Inches(2.3)
    tbl.columns[1].width = Inches(1.5)

    # Header row
    headers = ["Metric", "Q1 2025"]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.size = Pt(13)
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Data rows
    data_rows = [
        ("Total Revenue", "$84.2M"),
        ("EBITDA Margin", "34.5%"),
        ("Free Cash Flow", "$28.6M"),
        ("YoY Growth", "18.0%"),
    ]
    for r, (metric, value) in enumerate(data_rows, 1):
        for c, txt in enumerate([metric, value]):
            cell = tbl.cell(r, c)
            cell.text = txt
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
            # Alternating row fill
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF5)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # ---- Slide 4: Portfolio Highlights ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    fill4 = slide4.background.fill
    fill4.solid()
    fill4.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)

    add_title_box(slide4, "Portfolio Highlights",
                  Inches(0.5), Inches(0.4), Inches(9), Inches(0.7),
                  font_size=28, bold=True, color=RGBColor(0x1F, 0x49, 0x7D))

    rect4 = slide4.shapes.add_shape(
        1,
        Inches(0.5), Inches(1.15), Inches(9), Emu(36000)
    )
    rect4.fill.solid()
    rect4.fill.fore_color.rgb = RGBColor(0x1F, 0x49, 0x7D)
    rect4.line.fill.background()

    content4 = (
        "Top performing portfolio companies in Q1 2025:\n\n"
        "  • NovaTech Solutions — SaaS platform, ARR $42M (+34% YoY), Series C\n"
        "  • BlueHaven Logistics — Supply chain automation, EBITDA $18.5M, profitable\n"
        "  • Crestline Health — Digital therapeutics, patients served 1.2M, pre-IPO\n"
        "  • Arrowood Energy — Renewable infrastructure, capacity 480MW, cash-flow positive\n\n"
        "Two portfolio companies (NovaTech and Crestline Health) are on the IPO track "
        "for H2 2025, representing a combined estimated exit valuation of $1.8B."
    )
    add_body_box(slide4, content4,
                 Inches(0.5), Inches(1.4), Inches(9), Inches(5.5),
                 font_size=14, color=RGBColor(0x22, 0x22, 0x22))

    # ---- Slide 5: Outlook & Next Steps ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    fill5 = slide5.background.fill
    fill5.solid()
    fill5.fore_color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    add_title_box(slide5, "Outlook & Next Steps",
                  Inches(0.5), Inches(0.5), Inches(9), Inches(0.8),
                  font_size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    add_body_box(slide5,
                 "Q2 2025 Strategic Priorities",
                 Inches(0.5), Inches(1.5), Inches(9), Inches(0.6),
                 font_size=18, color=RGBColor(0xC8, 0xD8, 0xEB))

    content5 = (
        "  1. Close $220M fund extension for growth equity vertical by April 30\n"
        "  2. Complete secondary transactions for two legacy positions ($95M aggregate)\n"
        "  3. Advance NovaTech IPO preparation — S-1 draft filing target June 2025\n"
        "  4. Expand real assets vertical with two infrastructure acquisitions ($150M total)\n"
        "  5. Implement ESG reporting framework across 100% of portfolio companies"
    )
    add_body_box(slide5, content5,
                 Inches(0.5), Inches(2.2), Inches(9), Inches(4.0),
                 font_size=14, color=RGBColor(0xE8, 0xEE, 0xF5))

    add_body_box(slide5, "Meridian Capital Partners | Confidential | March 2025",
                 Inches(0.5), Inches(6.8), Inches(9), Inches(0.4),
                 font_size=11, color=RGBColor(0xA0, 0xB4, 0xCC))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
