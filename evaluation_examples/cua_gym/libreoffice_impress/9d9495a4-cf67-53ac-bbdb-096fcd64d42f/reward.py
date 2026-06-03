"""
Reward Script: Draw a connector (arrow) from the 'Input' box to the 'Process' box on slide 3.
Task ID: impress_objects_029
Domain: libreoffice_impress
Scoring:
  Component 1: A connector shape (LINE/cxnSp) exists on slide 3 (0.5 pts)
  Component 2: The connector has an arrowhead on at least one end (0.3 pts)
  Component 3: The connector is spatially positioned between the Input and Process boxes (0.2 pts)
Total: 1.0
"""

import os
import zipfile
import xml.etree.ElementTree as ET

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError as e:
    print(f"CRITICAL: Cannot import python-pptx: {e}")
    print("REWARD: 0.0")
    exit(0)

WORKDIR = '/home/user'
TASK_ID = 'impress_objects_029'

# XML namespaces for PPTX parsing
NS = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

# Connector geometry types recognized as line/connector shapes
CONNECTOR_PRST_TYPES = {
    'line', 'lineInv', 'straightConnector1',
    'bentConnector2', 'bentConnector3', 'bentConnector4', 'bentConnector5',
    'curvedConnector2', 'curvedConnector3', 'curvedConnector4', 'curvedConnector5',
}


def count_connectors_in_slide(slide, slide3_root):
    """
    Count connector shapes on a slide via both python-pptx API and raw XML parsing.
    Returns (api_count, cxnsp_count, line_sp_elements).
    """
    # Method A: python-pptx API — shape_type == MSO_SHAPE_TYPE.LINE
    line_shapes_api = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.LINE]

    # Method B: XML cxnSp elements (dedicated connector XML tag)
    cxnSp_elements = []
    if slide3_root is not None:
        cxnSp_elements = slide3_root.findall('.//p:cxnSp', NS)

    # Method C: p:sp with connector prstGeom type
    line_sp_elements = []
    if slide3_root is not None:
        for sp in slide3_root.findall('.//p:sp', NS):
            spPr = sp.find('p:spPr', NS)
            if spPr is not None:
                prstGeom = spPr.find('a:prstGeom', NS)
                if prstGeom is not None and prstGeom.get('prst') in CONNECTOR_PRST_TYPES:
                    line_sp_elements.append(sp)

    return line_shapes_api, cxnSp_elements, line_sp_elements


def check_arrow_in_connectors(cxnSp_elements, line_sp_elements):
    """
    Check if any connector has an arrowhead (headEnd or tailEnd with arrow type).
    Returns (found: bool, tag_name: str, arrow_type: str).
    """
    arrow_types = ('arrow', 'openArrow', 'stealth', 'triangle')
    all_connector_xml = list(cxnSp_elements) + list(line_sp_elements)
    for elem in all_connector_xml:
        for end_tag in ['a:headEnd', 'a:tailEnd']:
            end_elem = elem.find(f'.//{end_tag}', NS)
            if end_elem is not None:
                end_type = end_elem.get('type', 'none')
                if end_type in arrow_types:
                    return True, end_tag.split(':')[1], end_type
    return False, '', ''


def verify_task(file_path):
    """
    Verify that a connector (arrow) has been added from the 'Input' box
    to the 'Process' box on slide 3.

    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: file must be loadable
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, expected at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[2]  # Slide 3 (0-indexed)

    # Parse slide 3 XML once (used by all components)
    slide3_root = None
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide3.xml') as f:
                slide3_root = ET.fromstring(f.read().decode('utf-8'))
    except Exception as e:
        print(f"WARNING: Could not parse slide3.xml: {e}")

    # Gather connector shape data once for all components
    try:
        line_shapes_api, cxnSp_elements, line_sp_elements = count_connectors_in_slide(slide, slide3_root)
    except Exception as e:
        print(f"ERROR: Could not enumerate connector shapes: {e}")
        line_shapes_api, cxnSp_elements, line_sp_elements = [], [], []

    # -----------------------------------------------------------------------
    # Component 1: A connector shape exists on slide 3 (0.5 points)
    # -----------------------------------------------------------------------
    # The task adds a new connector between the boxes. Check for its presence via:
    # (a) python-pptx: shape.shape_type == MSO_SHAPE_TYPE.LINE
    # (b) XML: p:cxnSp element (dedicated connector XML tag)
    # (c) XML: p:sp with connector-type prstGeom
    # This fails on the initial file (no connector) and passes on the golden file.
    try:
        n_api = len(line_shapes_api)
        n_cxn = len(cxnSp_elements)
        n_sp = len(line_sp_elements)
        connector_found = (n_api > 0 or n_cxn > 0 or n_sp > 0)
        if connector_found:
            print(f"PASS: Component 1 — Connector shape(s) found on slide 3 "
                  f"(LINE via API: {n_api}, cxnSp elements: {n_cxn}, line sp: {n_sp}) (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 1 — No connector shape found on slide 3 "
                  f"(LINE via API: 0, cxnSp: 0, line sp: 0)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # -----------------------------------------------------------------------
    # Component 2: The connector has an arrowhead on at least one end (0.3 points)
    # -----------------------------------------------------------------------
    # The task says "connector (arrow)" — the connector must have an arrowhead.
    # In PPTX XML, arrowheads appear as <a:headEnd type="arrow"/> or <a:tailEnd type="arrow"/>
    # inside the <a:ln> element of the connector shape.
    try:
        has_arrow, arrow_end, arrow_type = check_arrow_in_connectors(cxnSp_elements, line_sp_elements)
        if has_arrow:
            print(f"PASS: Component 2 — Arrowhead found: {arrow_end} type='{arrow_type}' (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 — No arrowhead found on connector shapes "
                  f"(headEnd/tailEnd with arrow type missing)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # -----------------------------------------------------------------------
    # Component 3: Connector is positioned between the Input and Process boxes (0.2 points)
    # -----------------------------------------------------------------------
    # The connector should span the horizontal gap between the right edge of 'Input'
    # and the left edge of 'Process'.
    # Input box: left=914400, width=2286000 → right edge at 3200400 EMU
    # Process box: left=5943600 EMU
    # The connector must overlap with the horizontal range [3200400, 5943600].
    try:
        input_right = None
        process_left = None
        for shape in slide.shapes:
            if shape.name == 'Input':
                input_right = shape.left + shape.width
            elif shape.name == 'Process':
                process_left = shape.left

        if input_right is None or process_left is None:
            print(f"FAIL: Component 3 — Could not find 'Input' and/or 'Process' shapes "
                  f"(input_right={input_right}, process_left={process_left})")
        else:
            gap_left = input_right    # right edge of Input box
            gap_right = process_left  # left edge of Process box
            gap_width = gap_right - gap_left
            tolerance = max(gap_width * 0.25, 914400)  # at least 1 inch tolerance

            connector_in_gap = False

            # Check LINE shapes via python-pptx API
            for shape in line_shapes_api:
                conn_left = shape.left
                conn_right = shape.left + shape.width
                # Connector overlaps the gap region if ranges intersect (with tolerance)
                if conn_left <= (gap_right + tolerance) and conn_right >= (gap_left - tolerance):
                    connector_in_gap = True
                    print(f"PASS: Component 3 — Connector in gap between boxes "
                          f"(conn=[{conn_left}, {conn_right}], gap=[{gap_left}, {gap_right}]) (0.2 pts)")
                    break

            # Fall back to XML cxnSp position check
            if not connector_in_gap:
                for cxn in cxnSp_elements:
                    spPr = cxn.find('p:spPr', NS)
                    if spPr is not None:
                        xfrm = spPr.find('a:xfrm', NS)
                        if xfrm is not None:
                            off = xfrm.find('a:off', NS)
                            ext = xfrm.find('a:ext', NS)
                            if off is not None and ext is not None:
                                cx = int(off.get('x', 0))
                                cw = int(ext.get('cx', 0))
                                conn_left = cx
                                conn_right = cx + cw
                                if conn_left <= (gap_right + tolerance) and conn_right >= (gap_left - tolerance):
                                    connector_in_gap = True
                                    print(f"PASS: Component 3 — Connector (cxnSp) in gap "
                                          f"(conn=[{conn_left}, {conn_right}], gap=[{gap_left}, {gap_right}]) (0.2 pts)")
                                    break
                                else:
                                    print(f"FAIL: Component 3 — cxnSp not in gap "
                                          f"(conn=[{conn_left}, {conn_right}], gap=[{gap_left}, {gap_right}])")

            if not connector_in_gap and (len(line_shapes_api) + len(cxnSp_elements)) == 0:
                print(f"FAIL: Component 3 — No connector shape found to check position")
            elif not connector_in_gap:
                print(f"FAIL: Component 3 — Connector exists but is not in gap between Input and Process boxes")

            if connector_in_gap:
                total_score += 0.2
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against golden file (path on VM)
file_path = f'{WORKDIR}/{TASK_ID}_initial.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
