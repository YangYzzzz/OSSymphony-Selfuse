"""
Initial Setup: Draw a connector (arrow) from the 'Input' box to the 'Process' box on slide 3.
Task ID: impress_objects_029
Domain: libreoffice_impress

Creates flowchart_basic.pptx with 3 slides. Slide 3 has two rectangles:
'Input' on the left and 'Process' on the right. No connector exists yet.
"""

import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_objects_029'
OUTPUT = f'{WORKDIR}/{TASK_ID}_initial.pptx'
DESKTOP_PATH = f'{WORKDIR}/Desktop/flowchart_basic.pptx'


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Business Process Flowchart"
    slide1.placeholders[1].text = "Overview of Core Workflows\nQ1 2025"

    # Set title font
    title_tf = slide1.shapes.title.text_frame
    for para in title_tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(36)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Process Overview"
    content_tf = slide2.placeholders[1].text_frame
    content_tf.text = "This presentation covers the following workflows:"
    steps = [
        "Data Input Collection",
        "Data Processing and Transformation",
        "Output and Reporting",
    ]
    for step in steps:
        p = content_tf.add_paragraph()
        p.text = step
        p.level = 1
        for run in p.runs:
            run.font.size = Pt(20)

    # --- Slide 3: Flowchart Diagram (Input -> Process) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add a title text box at the top
    title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "Data Flow Diagram"
    p_title.alignment = PP_ALIGN.CENTER
    for run in p_title.runs:
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Add 'Input' rectangle on the left
    # Position: left=1.0", top=3.0", width=2.5", height=1.2"
    input_shape = slide3.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        Inches(1.0), Inches(3.0), Inches(2.5), Inches(1.2)
    )
    input_shape.name = "Input"
    input_fill = input_shape.fill
    input_fill.solid()
    input_fill.fore_color.rgb = RGBColor(0x4A, 0x86, 0xC8)

    input_tf = input_shape.text_frame
    input_tf.text = "Input"
    input_para = input_tf.paragraphs[0]
    input_para.alignment = PP_ALIGN.CENTER
    for run in input_para.runs:
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Add 'Process' rectangle on the right
    # Position: left=6.5", top=3.0", width=2.5", height=1.2"
    process_shape = slide3.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        Inches(6.5), Inches(3.0), Inches(2.5), Inches(1.2)
    )
    process_shape.name = "Process"
    process_fill = process_shape.fill
    process_fill.solid()
    process_fill.fore_color.rgb = RGBColor(0x70, 0xAD, 0x47)

    process_tf = process_shape.text_frame
    process_tf.text = "Process"
    process_para = process_tf.paragraphs[0]
    process_para.alignment = PP_ALIGN.CENTER
    for run in process_para.runs:
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Add a descriptive label below the shapes (no connector)
    label_box = slide3.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(8.0), Inches(0.6))
    label_tf = label_box.text_frame
    label_p = label_tf.paragraphs[0]
    label_p.text = "Note: Connect the Input and Process boxes with an arrow connector."
    label_p.alignment = PP_ALIGN.CENTER
    for run in label_p.runs:
        run.font.size = Pt(14)
        run.font.italic = True
        run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

    # Save the initial file
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Also copy to Desktop as flowchart_basic.pptx for the task context
    os.makedirs(f'{WORKDIR}/Desktop', exist_ok=True)
    shutil.copy(OUTPUT, DESKTOP_PATH)
    print(f'Copied to Desktop: {DESKTOP_PATH}')


create_initial()
