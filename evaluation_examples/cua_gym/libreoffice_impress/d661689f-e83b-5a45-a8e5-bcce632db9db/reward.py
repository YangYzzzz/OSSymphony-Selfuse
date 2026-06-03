"""
FINAL REWARD SCRIPT - SUCCESS
Task: In LibreOffice Impress, I nudged Picture 1 out of position on slide 179. Could you line it back up so it’s perfectly centered horizontally (vertical center) but still touching the bottom edge of the slide?
Generated: 2025-09-10 18:41:42
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation


def verify_picture_alignment(
    file_path: str,
    slide_index: int = 178,
    tolerance_ratio: float = 0.01,
    tolerance_emus: int = 10_000,
):
    """Verify that the (first) picture on the specified slide is
    1) perfectly centred horizontally and
    2) its bottom edge touches the bottom edge of the slide.

    Scoring (progressive):
        +0.7  if horizontal centre is within tolerance
        +0.3  if bottom edge is within tolerance
    Returns a float score in the range [0.0, 1.0].
    """

    # ------------------------------------------------------------------
    # 0.  Preliminary checks (file existence and slide availability)
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Unable to load presentation: {e}")
        return 0.0

    slide_count = len(prs.slides)
    print(f"✓ Loaded presentation – {slide_count} slides total")

    if slide_index >= slide_count:
        print(
            f"✗ Requested slide index {slide_index} (1-based {slide_index+1}) "
            f"but presentation only has {slide_count} slides"
        )
        return 0.0

    slide = prs.slides[slide_index]
    slide_w, slide_h = prs.slide_width, prs.slide_height
    print(f"Slide size (EMU): width={slide_w}, height={slide_h}")

    # ------------------------------------------------------------------
    # 1.  Locate the target picture shape
    # ------------------------------------------------------------------
    # In python-pptx the constant for pictures is 13 (MSO_SHAPE_TYPE.PICTURE)
    pic_shapes = [sh for sh in slide.shapes if sh.shape_type == 13]
    print(f"Found {len(pic_shapes)} picture shape(s) on slide {slide_index + 1}")

    if not pic_shapes:
        print("✗ No picture shapes found – cannot verify task")
        return 0.0

    # Prefer a shape whose name starts with "Picture 1", otherwise first picture
    target = None
    for sh in pic_shapes:
        if sh.name.lower().startswith("picture 1"):
            target = sh
            break
    if target is None:
        target = pic_shapes[0]
        print(f"✓ Using first picture named '{target.name}' for verification")
    else:
        print(f"✓ Using picture named '{target.name}' for verification")

    # ------------------------------------------------------------------
    # 2.  Geometric calculations
    # ------------------------------------------------------------------
    centre_x = target.left + target.width / 2
    expected_centre_x = slide_w / 2
    tol_x = max(slide_w * tolerance_ratio, tolerance_emus)

    bottom_y = target.top + target.height
    expected_bottom_y = slide_h
    tol_y = max(slide_h * tolerance_ratio, tolerance_emus)

    # ------------------------------------------------------------------
    # 3.  Verification and progressive scoring
    # ------------------------------------------------------------------
    score = 0.0

    # 3a. Horizontal centring check
    horiz_diff = abs(centre_x - expected_centre_x)
    horiz_ok = horiz_diff <= tol_x
    print(
        f"Horizontal centre difference: {horiz_diff} EMU (tolerance {tol_x}) – "
        f"{'OK' if horiz_ok else 'NOT OK'}"
    )
    if horiz_ok:
        score += 0.7

    # 3b. Bottom-edge alignment check
    bottom_diff = abs(bottom_y - expected_bottom_y)
    bottom_ok = bottom_diff <= tol_y
    print(
        f"Bottom edge difference: {bottom_diff} EMU (tolerance {tol_y}) – "
        f"{'OK' if bottom_ok else 'NOT OK'}"
    )
    if bottom_ok:
        score += 0.3

    # ------------------------------------------------------------------
    # 4.  Final score & report
    # ------------------------------------------------------------------
    final_score = min(score, 1.0)
    print(f"Final score: {final_score}")
    return final_score


if __name__ == "__main__":
    # Path provided by the task context
    FILE_PATH = (
        "/home/user/"
        "in_libreoffice_impress_i_nudged_picture_1_out_of_position_on_slide_"
        "179_could_you_line_it_back_up_so__golden.pptx"
    )

    reward = verify_picture_alignment(FILE_PATH)
    print(f"REWARD: {reward}")

