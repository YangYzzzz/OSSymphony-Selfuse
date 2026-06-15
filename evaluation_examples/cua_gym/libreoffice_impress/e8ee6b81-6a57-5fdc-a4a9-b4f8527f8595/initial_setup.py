"""
Initial Setup: Create marketing_deck.pptx with 5 slides. Slide 5 'Development Process' is empty.
Task ID: impress_gf5_011
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_011'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_placeholder(placeholder, text, font_name="Arial", font_size=Pt(18),
                            bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to set text in a placeholder with formatting."""
    tf = placeholder.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    for run in p.runs:
        run.font.name = font_name
        run.font.size = font_size
        run.font.bold = bold
        if color:
            run.font.color.rgb = color


def add_bullet_slide(prs, title_text, bullets):
    """Add a slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide.shapes.title.text = title_text
    # Set title font
    for run in slide.shapes.title.text_frame.paragraphs[0].runs:
        run.font.name = "Arial"
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # Add bullets to the content placeholder
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet_text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet_text
        p.level = 0
        p.space_after = Pt(6)
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)
    return slide


def create_initial():
    prs = Presentation()
    # Standard 16:9 dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    slide1.shapes.title.text = "Marketing Deck"
    for run in slide1.shapes.title.text_frame.paragraphs[0].runs:
        run.font.name = "Arial"
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    slide1.placeholders[1].text = "Q4 2025 Investor Strategy Presentation"
    for run in slide1.placeholders[1].text_frame.paragraphs[0].runs:
        run.font.name = "Arial"
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)

    # --- Slide 2: Market Overview ---
    add_bullet_slide(prs, "Market Overview", [
        "Global SaaS market projected to reach $908B by 2030",
        "Enterprise adoption rate increased 34% year-over-year",
        "Key verticals: Healthcare, Finance, Education, Retail",
        "Competitive landscape consolidating around 5 major players",
        "Our market share grew from 4.2% to 6.8% in 2025",
    ])

    # --- Slide 3: Target Audience ---
    add_bullet_slide(prs, "Target Audience", [
        "Mid-market enterprises (500-5000 employees)",
        "CTO and VP Engineering decision makers",
        "Organizations with hybrid cloud infrastructure",
        "Companies spending $500K+ annually on developer tools",
        "Industries with strict compliance requirements (SOC2, HIPAA)",
    ])

    # --- Slide 4: Campaign Strategy ---
    add_bullet_slide(prs, "Campaign Strategy", [
        "Launch multi-channel digital campaign in October 2025",
        "Partner with 12 industry influencers for co-branded content",
        "Host 3 regional customer summits (NYC, London, Singapore)",
        "Allocate $2.4M budget across paid, organic, and events",
        "Target 15,000 qualified leads by end of Q4",
    ])

    # --- Slide 5: Development Process (EMPTY — title only) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only layout
    slide5.shapes.title.text = "Development Process"
    for run in slide5.shapes.title.text_frame.paragraphs[0].runs:
        run.font.name = "Arial"
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
