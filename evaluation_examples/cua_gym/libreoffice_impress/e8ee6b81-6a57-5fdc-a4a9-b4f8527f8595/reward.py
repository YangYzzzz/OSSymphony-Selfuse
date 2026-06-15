"""
Reward Script: SmartArt-style process diagram on slide 5
Task ID: impress_gf5_011
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): 4 rectangular shapes with correct labels on slide 5
  Component 2 (0.25): Arrow connectors between the 4 boxes (3 arrows)
  Component 3 (0.25): Uniform fill color across all 4 boxes
  Component 4 (0.20): Uniform size and vertical alignment of all 4 boxes
"""

import os
import sys

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_011'

# Persistence hook for LibreOffice Impress
def persist_app_state():
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        import time
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # 0-indexed, slide 5

    # Collect all shapes on slide 5, categorize them
    EXPECTED_LABELS = ['Research', 'Design', 'Develop', 'Deploy']
    rect_shapes = []  # shapes with expected label text
    arrow_shapes = []  # shapes that look like arrows (no meaningful text, or arrow in name)

    for shape in slide.shapes:
        # Skip the title placeholder
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            continue

        shape_text = ""
        if hasattr(shape, 'text'):
            shape_text = shape.text.strip()

        # Classify: if text matches one of our labels, it's a rect box
        if shape_text in EXPECTED_LABELS:
            rect_shapes.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # Could be an arrow or connector shape
            arrow_shapes.append(shape)

    print(f"INFO: Found {len(rect_shapes)} labeled rectangles, {len(arrow_shapes)} other auto-shapes (potential arrows)")

    # Component 1: 4 rectangular shapes with correct labels on slide 5 (0.30 points)
    try:
        found_labels = set()
        for shape in rect_shapes:
            label = shape.text.strip()
            found_labels.add(label)

        missing = set(EXPECTED_LABELS) - found_labels
        found_count = len(found_labels.intersection(EXPECTED_LABELS))

        if found_count == 4:
            print(f"PASS: Component 1 -- All 4 labels found: {sorted(found_labels)} (0.30 pts)")
            total_score += 0.30
        elif found_count >= 2:
            partial = round(0.30 * (found_count / 4), 2)
            print(f"PARTIAL: Component 1 -- {found_count}/4 labels found, missing: {missing} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- Only {found_count}/4 labels found. Missing: {missing}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Arrow connectors between the 4 boxes (0.25 points)
    # We expect 3 arrows for Research->Design->Develop->Deploy
    try:
        # Count shapes that are likely arrows: auto-shapes that are NOT labeled rectangles
        # Arrows typically have no text or very short text, and are smaller than the boxes
        arrow_count = len(arrow_shapes)

        if arrow_count >= 3:
            print(f"PASS: Component 2 -- {arrow_count} arrow/connector shapes found (need >= 3) (0.25 pts)")
            total_score += 0.25
        elif arrow_count >= 1:
            partial = round(0.25 * (arrow_count / 3), 2)
            print(f"PARTIAL: Component 2 -- {arrow_count}/3 arrows found ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 -- No arrow/connector shapes found between boxes")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Uniform fill color across all 4 boxes (0.25 points)
    try:
        if len(rect_shapes) >= 4:
            fill_colors = []
            for shape in rect_shapes:
                try:
                    fill = shape.fill
                    if fill.type == 1:  # SOLID
                        fill_colors.append(str(fill.fore_color.rgb))
                    else:
                        fill_colors.append(f"type_{fill.type}")
                except Exception:
                    fill_colors.append("unknown")

            unique_colors = set(fill_colors)
            has_solid_fill = all(c not in ("unknown", "type_None") for c in fill_colors)

            if len(unique_colors) == 1 and has_solid_fill:
                print(f"PASS: Component 3 -- All 4 boxes have uniform fill color: {fill_colors[0]} (0.25 pts)")
                total_score += 0.25
            elif has_solid_fill and len(unique_colors) <= 2:
                print(f"PARTIAL: Component 3 -- Boxes have fills but not uniform: {fill_colors} (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 3 -- Fill colors not uniform or missing: {fill_colors}")
        else:
            print(f"FAIL: Component 3 -- Need 4 labeled boxes to check fill, found {len(rect_shapes)}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Uniform size and vertical alignment of all 4 boxes (0.20 points)
    try:
        if len(rect_shapes) >= 4:
            widths = [s.width for s in rect_shapes]
            heights = [s.height for s in rect_shapes]
            tops = [s.top for s in rect_shapes]

            # Check uniform size (all widths same, all heights same)
            # Use 5% tolerance
            def all_close(values, tol=0.05):
                if not values:
                    return False
                ref = values[0]
                if ref == 0:
                    return all(v == 0 for v in values)
                return all(abs(v - ref) / ref <= tol for v in values)

            size_uniform = all_close(widths) and all_close(heights)
            # Check vertical alignment (all tops are approximately same)
            top_aligned = all_close(tops)

            if size_uniform and top_aligned:
                print(f"PASS: Component 4 -- Uniform size ({widths[0]}x{heights[0]}) and aligned tops ({tops[0]}) (0.20 pts)")
                total_score += 0.20
            elif size_uniform:
                print(f"PARTIAL: Component 4 -- Uniform size but tops not aligned: {tops} (0.10 pts)")
                total_score += 0.10
            elif top_aligned:
                print(f"PARTIAL: Component 4 -- Tops aligned but sizes differ: w={widths}, h={heights} (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 4 -- Sizes differ (w={widths}, h={heights}) and tops not aligned ({tops})")
        else:
            print(f"FAIL: Component 4 -- Need 4 labeled boxes, found {len(rect_shapes)}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
