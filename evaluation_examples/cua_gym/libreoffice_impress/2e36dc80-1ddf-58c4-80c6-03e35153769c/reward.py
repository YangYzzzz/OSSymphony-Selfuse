"""
Reward Script: Resize images on slides 3-7 to 4 inches width, maintain aspect ratio, center horizontally
Task ID: impress_stu_040
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): All 5 images on slides 3-7 have width = 4in AND correct height (aspect ratio preserved)
  Component 2 (0.5): All 5 images are horizontally centered on their slides
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_040'
TARGET_WIDTH_EMU = 3657600  # 4 inches = 4 * 914400 EMU
WIDTH_TOLERANCE = 0.02  # 2% tolerance
CENTER_TOLERANCE = 0.03  # 3% tolerance

# Expected heights after resizing to 4in width while preserving original aspect ratios
# Calculated as: original_height / original_width * 4_inches_in_emu
# These are the EXPECTED golden heights; they only pass if width was changed to 4in AND ratio preserved
EXPECTED_HEIGHTS = {
    2: int(round(3771900 / 5029200 * TARGET_WIDTH_EMU)),   # Slide 3: 2743200
    3: int(round(4114800 / 2743200 * TARGET_WIDTH_EMU)),   # Slide 4: 5486400
    4: int(round(3200400 / 6400800 * TARGET_WIDTH_EMU)),   # Slide 5: 1828800
    5: int(round(2286000 / 2286000 * TARGET_WIDTH_EMU)),   # Slide 6: 3657600
    6: int(round(2438400 / 5486400 * TARGET_WIDTH_EMU)),   # Slide 7: 1625600 (approx)
}

# Slides to check (0-indexed)
TARGET_SLIDES = [2, 3, 4, 5, 6]


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    if num_slides < 7:
        print(f"CRITICAL: Expected at least 7 slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    slide_width_emu = prs.slide_width

    # Collect image data from target slides
    image_data = {}
    for idx in TARGET_SLIDES:
        slide = prs.slides[idx]
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_data[idx] = {
                    'width': shape.width,
                    'height': shape.height,
                    'left': shape.left,
                }
                break  # One image per slide expected

    if len(image_data) < 5:
        print(f"CRITICAL: Expected images on slides 3-7, found {len(image_data)} images on target slides")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Width = 4 inches AND correct height (aspect ratio preserved) (0.5 points)
    # Per-image: 0.10 points each
    # This component ONLY passes when width has been changed to 4in AND height adjusted proportionally.
    # On the initial file, widths vary (5.5, 3.0, 7.0, 2.5, 6.0 inches) so this will FAIL.
    try:
        comp1_pass = 0
        comp1_score = 0.0
        for idx in TARGET_SLIDES:
            img = image_data[idx]
            actual_width = img['width']
            actual_height = img['height']
            expected_height = EXPECTED_HEIGHTS[idx]

            width_ok = abs(actual_width - TARGET_WIDTH_EMU) / TARGET_WIDTH_EMU <= WIDTH_TOLERANCE
            height_ok = abs(actual_height - expected_height) / max(expected_height, 1) <= WIDTH_TOLERANCE

            if width_ok and height_ok:
                comp1_pass += 1
                comp1_score += 0.10
                print(f"PASS: Slide {idx+1} size = {actual_width/914400:.2f}in x {actual_height/914400:.2f}in (expected 4.00in x {expected_height/914400:.2f}in)")
            else:
                reasons = []
                if not width_ok:
                    reasons.append(f"width={actual_width/914400:.2f}in (expected 4.00in)")
                if not height_ok:
                    reasons.append(f"height={actual_height/914400:.2f}in (expected {expected_height/914400:.2f}in)")
                print(f"FAIL: Slide {idx+1} -- {', '.join(reasons)}")
        comp1_score = min(comp1_score, 0.5)
        if comp1_score > 0:
            total_score += comp1_score
        print(f"Component 1 (width+ratio): {comp1_pass}/5 passed, {comp1_score:.2f}/0.50 pts")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Images horizontally centered (0.5 points)
    # Per-image: 0.10 points each
    # Centered means: left = (slide_width - image_width) / 2
    # On the initial file, images are at various non-centered positions, so this will FAIL.
    try:
        comp2_pass = 0
        comp2_score = 0.0
        for idx in TARGET_SLIDES:
            img = image_data[idx]
            expected_left = (slide_width_emu - img['width']) / 2
            actual_left = img['left']
            if expected_left > 0:
                rel_diff = abs(actual_left - expected_left) / expected_left
            else:
                # If image is wider than slide, expected_left could be 0 or negative
                rel_diff = abs(actual_left - expected_left) / slide_width_emu if slide_width_emu > 0 else 0
            if rel_diff <= CENTER_TOLERANCE:
                comp2_pass += 1
                comp2_score += 0.10
                print(f"PASS: Slide {idx+1} left = {actual_left/914400:.2f}in, expected {expected_left/914400:.2f}in (diff {rel_diff*100:.1f}%)")
            else:
                print(f"FAIL: Slide {idx+1} left = {actual_left/914400:.2f}in, expected {expected_left/914400:.2f}in (diff {rel_diff*100:.1f}%)")
        comp2_score = min(comp2_score, 0.5)
        if comp2_score > 0:
            total_score += comp2_score
        print(f"Component 2 (centered): {comp2_pass}/5 passed, {comp2_score:.2f}/0.50 pts")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
