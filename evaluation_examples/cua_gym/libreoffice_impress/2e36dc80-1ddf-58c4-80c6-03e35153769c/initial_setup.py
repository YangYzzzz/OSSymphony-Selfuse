"""
Initial Setup: Resize and center images on slides 3-7
Task ID: impress_stu_040
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import io

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image as PILImage

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_040'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def make_dummy_image(width_px, height_px, color, label=""):
    """Create an in-memory PNG image for embedding into the presentation."""
    img = PILImage.new('RGB', (width_px, height_px), color)
    buf = io.BytesIO()
    img.save(buf, format='PNG')
    buf.seek(0)
    return buf


def create_initial():
    prs = Presentation()
    # Standard slide dimensions: 10x7.5 inches (landscape)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Urban Architecture: A Photo Essay"
    slide1.placeholders[1].text = "Exploring Modern City Landscapes Through Photography"

    # --- Slide 2: Introduction ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide2.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Introduction"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x62)

    body = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    btf = body.text_frame
    btf.word_wrap = True
    bp = btf.paragraphs[0]
    bp.text = (
        "This photo essay captures the intersection of historical and contemporary "
        "architectural styles found across major metropolitan areas. Each image "
        "highlights a unique aspect of how cities evolve while preserving their heritage."
    )
    bp.runs[0].font.size = Pt(16)

    # --- Slides 3-7: Photo slides with images of varying sizes and positions ---
    # Each slide has a title and one image at different sizes/positions
    photo_data = [
        {
            "title": "The Glass Tower - Downtown Financial District",
            "img_w_px": 800, "img_h_px": 600, "color": (70, 130, 180),
            "pptx_w": Inches(5.5), "pptx_left": Inches(0.5), "pptx_top": Inches(1.8),
        },
        {
            "title": "Historic Brownstone Row - Beacon Hill",
            "img_w_px": 600, "img_h_px": 900, "color": (139, 90, 43),
            "pptx_w": Inches(3.0), "pptx_left": Inches(6.0), "pptx_top": Inches(1.5),
        },
        {
            "title": "Waterfront Promenade at Sunset",
            "img_w_px": 1000, "img_h_px": 500, "color": (255, 140, 0),
            "pptx_w": Inches(7.0), "pptx_left": Inches(0.3), "pptx_top": Inches(2.0),
        },
        {
            "title": "The Spiral Museum - Contemporary Arts Quarter",
            "img_w_px": 700, "img_h_px": 700, "color": (100, 100, 100),
            "pptx_w": Inches(2.5), "pptx_left": Inches(1.0), "pptx_top": Inches(2.5),
        },
        {
            "title": "Railway Station Canopy - Industrial Heritage",
            "img_w_px": 900, "img_h_px": 400, "color": (85, 107, 47),
            "pptx_w": Inches(6.0), "pptx_left": Inches(3.5), "pptx_top": Inches(1.2),
        },
    ]

    for i, pd in enumerate(photo_data):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank

        # Add title text
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = pd["title"]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        # Create image and add at non-standard size and position (NOT centered, NOT 4 inches)
        img_buf = make_dummy_image(pd["img_w_px"], pd["img_h_px"], pd["color"])
        # Compute height from aspect ratio
        aspect = pd["img_h_px"] / pd["img_w_px"]
        pptx_h = int(pd["pptx_w"] * aspect)
        pic = slide.shapes.add_picture(
            img_buf, pd["pptx_left"], pd["pptx_top"], pd["pptx_w"], pptx_h
        )

    # --- Slide 8: Reflections ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide8.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Reflections"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x62)

    body8 = slide8.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    btf8 = body8.text_frame
    btf8.word_wrap = True
    bp8 = btf8.paragraphs[0]
    bp8.text = (
        "Architecture tells the story of a city's aspirations. From towering glass facades "
        "to preserved brownstone rows, each structure carries the weight of decisions made "
        "by planners, architects, and communities over decades."
    )
    bp8.runs[0].font.size = Pt(16)

    # --- Slide 9: Credits ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide9.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Photography & Essay by Elena Vasquez"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    p2 = tf.add_paragraph()
    p2.text = "March 2026"
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = ""
    p2.runs[0].font.size = Pt(18)
    p2.runs[0].font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the presentation
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
