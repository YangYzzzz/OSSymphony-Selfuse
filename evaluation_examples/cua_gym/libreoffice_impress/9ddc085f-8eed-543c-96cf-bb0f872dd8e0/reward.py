"""
Reward Script: Insert doughnut chart on slide 5 with grade distribution
Task ID: impress_teach_070
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): Slide 5 contains a chart shape
  Component 2 (0.15): Chart is doughnut type
  Component 3 (0.15): Chart title is 'Final Grade Distribution'
  Component 4 (0.25): Chart categories and values match specification
  Component 5 (0.25): Chart data point colors match specification
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_070'

# Expected chart data
EXPECTED_CATEGORIES = ['A', 'B', 'C', 'D', 'F']
EXPECTED_VALUES = [20.0, 35.0, 25.0, 12.0, 8.0]
EXPECTED_COLORS = {
    0: '4CAF50',  # A = green
    1: '2196F3',  # B = blue
    2: 'FFC107',  # C = yellow
    3: 'FF9800',  # D = orange
    4: 'F44336',  # F = red
}
EXPECTED_TITLE = 'Final Grade Distribution'


def persist_app_state(domain):
    """Save any unsaved LibreOffice edits before verification."""
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print(f"PERSIST: ctrl+s sent for {domain}")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide5 = prs.slides[4]  # 0-indexed

    # Find chart shape on slide 5
    chart_shape = None
    for shape in slide5.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Slide 5 contains a chart shape (0.20 points)
    try:
        if chart_shape is not None:
            print(f"PASS: Component 1 -- Slide 5 has a chart shape (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 -- No chart shape found on slide 5")
            # No chart means all other components will also fail
            print(f"\nScore: {total_score}/1.0")
            print(f"REWARD: {total_score}")
            return total_score
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    chart = chart_shape.chart

    # Component 2: Chart is doughnut type (0.15 points)
    try:
        from pptx.enum.chart import XL_CHART_TYPE
        ct = chart.chart_type
        # DOUGHNUT enum value is -4120
        if ct == XL_CHART_TYPE.DOUGHNUT:
            print(f"PASS: Component 2 -- Chart type is DOUGHNUT (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 -- Expected DOUGHNUT, found {ct}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Chart title is 'Final Grade Distribution' (0.15 points)
    try:
        if chart.has_title:
            actual_title = chart.chart_title.text_frame.text.strip()
            if actual_title == EXPECTED_TITLE:
                print(f"PASS: Component 3 -- Chart title is '{actual_title}' (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 3 -- Expected title '{EXPECTED_TITLE}', found '{actual_title}'")
        else:
            print(f"FAIL: Component 3 -- Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Chart categories and values match (0.25 points)
    try:
        series = chart.series[0]
        actual_values = list(series.values)
        actual_categories = [str(c) for c in chart.plots[0].categories]

        cats_match = actual_categories == EXPECTED_CATEGORIES
        vals_match = len(actual_values) == len(EXPECTED_VALUES) and all(
            abs(a - e) < 0.01 for a, e in zip(actual_values, EXPECTED_VALUES)
        )

        if cats_match and vals_match:
            print(f"PASS: Component 4 -- Categories {actual_categories} and values {actual_values} match (0.25 pts)")
            total_score += 0.25
        else:
            if not cats_match:
                print(f"FAIL: Component 4 -- Categories mismatch: expected {EXPECTED_CATEGORIES}, found {actual_categories}")
            if not vals_match:
                print(f"FAIL: Component 4 -- Values mismatch: expected {EXPECTED_VALUES}, found {actual_values}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Chart data point colors match specification (0.25 points)
    try:
        series = chart.series[0]
        color_matches = 0
        total_points = len(EXPECTED_COLORS)

        for pt_idx, expected_hex in EXPECTED_COLORS.items():
            try:
                pt = series.points[pt_idx]
                fill = pt.format.fill
                if fill.type is not None:
                    actual_rgb = str(fill.fore_color.rgb).upper()
                    expected_upper = expected_hex.upper()
                    if actual_rgb == expected_upper:
                        color_matches += 1
                        print(f"  Color point {pt_idx}: MATCH ({actual_rgb})")
                    else:
                        print(f"  Color point {pt_idx}: MISMATCH expected {expected_upper}, found {actual_rgb}")
                else:
                    print(f"  Color point {pt_idx}: No solid fill set")
            except Exception as e:
                print(f"  Color point {pt_idx}: ERROR -- {e}")

        if color_matches == total_points:
            print(f"PASS: Component 5 -- All {total_points} colors match (0.25 pts)")
            total_score += 0.25
        elif color_matches > 0:
            partial = 0.25 * (color_matches / total_points)
            print(f"PARTIAL: Component 5 -- {color_matches}/{total_points} colors match ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 -- No colors match")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
