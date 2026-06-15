"""
Initial Setup: Create a 7-slide Semester Report presentation with Grade Overview on slide 5 (no chart).
Task ID: impress_teach_070
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_070'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_body(slide, title_text, body_lines):
    """Set title and body text on a slide with Title+Content layout."""
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0


def add_title_only(slide, title_text):
    """Set title on a title-only layout slide."""
    slide.shapes.title.text = title_text


def create_initial():
    prs = Presentation()
    # Standard widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    layout_title = prs.slide_layouts[0]       # Title Slide
    layout_content = prs.slide_layouts[1]     # Title + Content
    layout_title_only = prs.slide_layouts[5]  # Blank

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(layout_title)
    slide1.shapes.title.text = "Fall 2024 Semester Report"
    slide1.placeholders[1].text = "Department of Computer Science\nPrepared by Dr. Elena Martinez"

    # --- Slide 2: Course Enrollment ---
    slide2 = prs.slides.add_slide(layout_content)
    add_title_and_body(slide2, "Course Enrollment", [
        "Total students enrolled: 487 across 12 sections",
        "CS 101 Introduction to Programming: 142 students",
        "CS 201 Data Structures: 98 students",
        "CS 301 Algorithms: 76 students",
        "CS 401 Machine Learning: 63 students",
        "CS 450 Computer Vision: 54 students",
        "CS 499 Senior Capstone: 54 students",
    ])

    # --- Slide 3: Attendance Summary ---
    slide3 = prs.slides.add_slide(layout_content)
    add_title_and_body(slide3, "Attendance Summary", [
        "Average attendance rate: 84.2%",
        "Highest attendance: CS 101 (91.3%)",
        "Lowest attendance: CS 401 (72.8%)",
        "Students with perfect attendance: 38 (7.8%)",
        "Attendance improved 3.5% from Spring 2024",
        "Friday sessions had 12% lower attendance on average",
    ])

    # --- Slide 4: Assignment Completion ---
    slide4 = prs.slides.add_slide(layout_content)
    add_title_and_body(slide4, "Assignment Completion", [
        "Total assignments issued: 2,847 across all courses",
        "On-time submission rate: 78.4%",
        "Late submissions (within grace period): 14.2%",
        "Missing submissions: 7.4%",
        "Average score on completed assignments: 81.6/100",
        "Most challenging assignment: CS 301 HW5 (avg 62.3/100)",
    ])

    # --- Slide 5: Grade Overview (NO CHART - just title) ---
    slide5 = prs.slides.add_slide(layout_title_only)
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Grade Overview"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Add a subtitle text box below
    txBox2 = slide5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Distribution of final grades across all courses for Fall 2024"
    run2 = p2.runs[0]
    run2.font.size = Pt(18)
    run2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # --- Slide 6: Student Feedback ---
    slide6 = prs.slides.add_slide(layout_content)
    add_title_and_body(slide6, "Student Feedback", [
        "Overall satisfaction score: 4.2 / 5.0",
        "Course content relevance: 4.4 / 5.0",
        "Instructor effectiveness: 4.1 / 5.0",
        "Lab and resource availability: 3.8 / 5.0",
        "Top request: More hands-on projects (mentioned by 62% of students)",
        "219 students submitted end-of-semester evaluations",
    ])

    # --- Slide 7: Next Steps ---
    slide7 = prs.slides.add_slide(layout_content)
    add_title_and_body(slide7, "Next Steps", [
        "Increase TA support for CS 301 and CS 401",
        "Pilot flipped classroom model in CS 201 (Spring 2025)",
        "Expand office hours by 20% based on student feedback",
        "Introduce peer mentoring program for at-risk students",
        "Review and update CS 450 curriculum for industry alignment",
        "Target 88% attendance rate for Spring 2025 semester",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
