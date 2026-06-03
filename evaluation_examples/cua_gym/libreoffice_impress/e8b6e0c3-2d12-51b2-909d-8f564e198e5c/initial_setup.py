"""
Initial Setup: 7-slide consulting report presentation
Task ID: osworld_impress_title_color_match_003
Domain: libreoffice_impress

Creates a consulting report presentation where:
- Slide 3 title uses deep burgundy color (#722F37)
- Slide 5 title reads 'Action Items' in BLACK (not burgundy)
The task is to rename slide 5's title to 'Recommendations' and match its color to slide 3.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_color_match_003'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

BURGUNDY = RGBColor(0x72, 0x2F, 0x37)
BLACK    = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_title_text_color(slide, title_text, color):
    """Set title placeholder text and color on a slide."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.shape_type == 13:
            continue
        if shape.has_text_frame:
            # Check if it's the title (placeholder index 0 or named Title)
            try:
                if shape.placeholder_format is not None and shape.placeholder_format.idx == 0:
                    tf = shape.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = title_text
                    run.font.color.rgb = color
                    run.font.size = Pt(32)
                    run.font.bold = True
                    return
            except Exception:
                pass


def add_title_shape(slide, title_text, color, prs):
    """Add a title text box at the top of a slide."""
    left = Inches(0.5)
    top = Inches(0.3)
    width = prs.slide_width - Inches(1.0)
    height = Inches(1.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = color
    return txBox


def add_body_text(slide, text, top_offset, prs):
    """Add a body text box to a slide."""
    left = Inches(0.5)
    top = Inches(top_offset)
    width = prs.slide_width - Inches(1.0)
    height = Inches(5.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(18)
        run.font.color.rgb = DARK_GRAY
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Strategic Growth Initiative"
    slide1.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK_GRAY
    slide1.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(36)
    try:
        slide1.placeholders[1].text = "Q1 2025 Consulting Report\nMercer & Partners Advisory"
    except Exception:
        pass

    # ---- Slide 2: Executive Summary ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    slide2.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK_GRAY
    slide2.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    try:
        body2 = slide2.placeholders[1]
        tf2 = body2.text_frame
        tf2.text = "Market conditions present significant opportunity for expansion."
        tf2.add_paragraph().text = "Revenue targets exceeded by 12% in Q4 2024."
        tf2.add_paragraph().text = "Three strategic pillars identified for 2025 growth."
        tf2.add_paragraph().text = "Operational efficiency improvements yield $2.4M savings."
    except Exception:
        pass

    # ---- Slide 3: Market Analysis (title in BURGUNDY #722F37) ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Market Analysis"
    # Apply deep burgundy color to slide 3 title
    for para in slide3.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = BURGUNDY
            run.font.size = Pt(32)
    try:
        body3 = slide3.placeholders[1]
        tf3 = body3.text_frame
        tf3.text = "Total addressable market: $4.8B globally."
        tf3.add_paragraph().text = "Asia-Pacific segment growing at 18% CAGR."
        tf3.add_paragraph().text = "Key competitors: Nexus Corp, Hallmark Ventures, TerraGroup."
        tf3.add_paragraph().text = "Customer acquisition cost down 22% YoY."
        tf3.add_paragraph().text = "NPS score improved from 34 to 61 over 18 months."
    except Exception:
        pass

    # ---- Slide 4: Financial Overview ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Financial Overview"
    slide4.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK_GRAY
    slide4.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    try:
        body4 = slide4.placeholders[1]
        tf4 = body4.text_frame
        tf4.text = "FY2024 Revenue: $38.7M (+15% YoY)"
        tf4.add_paragraph().text = "EBITDA Margin: 24.3%"
        tf4.add_paragraph().text = "Operating Cash Flow: $9.2M"
        tf4.add_paragraph().text = "Debt-to-Equity Ratio: 0.42"
        tf4.add_paragraph().text = "Capital expenditure forecast: $5.1M for FY2025"
    except Exception:
        pass

    # ---- Slide 5: Action Items (title in BLACK — task requires changing to Recommendations + BURGUNDY) ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Action Items"
    # Title in BLACK (NOT burgundy — this is the initial state that needs to be changed)
    for para in slide5.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = BLACK
            run.font.size = Pt(32)
    try:
        body5 = slide5.placeholders[1]
        tf5 = body5.text_frame
        tf5.text = "Launch pilot program in Southeast Asian markets by June 2025."
        tf5.add_paragraph().text = "Hire 12 additional sales representatives across APAC."
        tf5.add_paragraph().text = "Renegotiate supplier contracts to achieve 8% cost reduction."
        tf5.add_paragraph().text = "Implement CRM upgrade to Salesforce Enterprise edition."
        tf5.add_paragraph().text = "Schedule quarterly board review for June 30, 2025."
    except Exception:
        pass

    # ---- Slide 6: Risk Assessment ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Risk Assessment"
    slide6.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK_GRAY
    slide6.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    try:
        body6 = slide6.placeholders[1]
        tf6 = body6.text_frame
        tf6.text = "Regulatory risk: Medium — new compliance requirements in EU by Q3."
        tf6.add_paragraph().text = "Currency risk: Low — 82% revenue in USD."
        tf6.add_paragraph().text = "Supply chain disruption: Medium — dual-sourcing strategy in progress."
        tf6.add_paragraph().text = "Talent retention: High — retention program budgeted at $1.2M."
    except Exception:
        pass

    # ---- Slide 7: Conclusion ----
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Conclusion"
    slide7.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK_GRAY
    slide7.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    try:
        body7 = slide7.placeholders[1]
        tf7 = body7.text_frame
        tf7.text = "The company is well-positioned for sustained growth in 2025."
        tf7.add_paragraph().text = "Execution of the three strategic pillars is critical to success."
        tf7.add_paragraph().text = "Next review scheduled: April 15, 2025."
    except Exception:
        pass

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
