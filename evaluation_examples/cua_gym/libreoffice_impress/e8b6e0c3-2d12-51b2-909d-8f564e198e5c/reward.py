"""
Reward Script: Rename slide 5's title to 'Recommendations' and match its color to the slide 3 title color.
Task ID: osworld_impress_title_color_match_003
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 5 title text == 'Recommendations'  — 0.5 points
  Component 2: Slide 5 title color == #722F37            — 0.5 points
  Total: 1.0
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_color_match_003'

# Expected ground-truth values (from task context)
EXPECTED_SLIDE5_TITLE_TEXT = 'Recommendations'
EXPECTED_SLIDE5_TITLE_COLOR = '722F37'  # deep burgundy, same as slide 3


def get_title_shape(slide):
    """Return the title placeholder shape from a slide, or None."""
    for shape in slide.shapes:
        if shape.has_text_frame and hasattr(shape, 'placeholder_format') and shape.placeholder_format:
            if shape.placeholder_format.idx == 0:
                return shape
    return None


def get_title_text(slide):
    """Return the full text of the title placeholder on the slide."""
    shape = get_title_shape(slide)
    if shape is None:
        return None
    return shape.text_frame.text.strip()


def get_title_color_rgb(slide):
    """
    Return the RGB hex string (uppercase) of the first run on the title
    placeholder, or None if no explicit color is set.
    Checks: run.font.color.type is not None → run.font.color.rgb.
    """
    shape = get_title_shape(slide)
    if shape is None:
        return None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if (run.text or '').strip():
                try:
                    if run.font.color.type is not None:
                        return str(run.font.color.rgb).upper()
                except Exception:
                    pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"CRITICAL: Presentation has only {len(prs.slides)} slides, expected at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide5 = prs.slides[4]  # 0-indexed → slide 5

    # -------------------------------------------------------------------
    # Component 1: Slide 5 title text is 'Recommendations' (0.5 points)
    # This FAILS on initial (text = 'Action Items') and PASSES on golden.
    # -------------------------------------------------------------------
    try:
        actual_text = get_title_text(slide5)
        if actual_text == EXPECTED_SLIDE5_TITLE_TEXT:
            print(f"PASS: Component 1 — Slide 5 title text is '{actual_text}' (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 1 — Slide 5 title text expected '{EXPECTED_SLIDE5_TITLE_TEXT}', found '{actual_text}'")
    except Exception as e:
        print(f"ERROR: Component 1 — Could not read slide 5 title text: {e}")

    # -------------------------------------------------------------------
    # Component 2: Slide 5 title color is #722F37 (0.5 points)
    # This FAILS on initial (color = #000000) and PASSES on golden.
    # -------------------------------------------------------------------
    try:
        actual_color = get_title_color_rgb(slide5)
        if actual_color == EXPECTED_SLIDE5_TITLE_COLOR.upper():
            print(f"PASS: Component 2 — Slide 5 title color is #{actual_color} (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 2 — Slide 5 title color expected #{EXPECTED_SLIDE5_TITLE_COLOR.upper()}, found #{actual_color}")
    except Exception as e:
        print(f"ERROR: Component 2 — Could not read slide 5 title color: {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the VM environment
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
