"""
Reward script for impress_teach_040.
Verifies that slide 7 has a line chart with correct data, title, and line color.
"""
import os
import time

# ── persist any unsaved LibreOffice changes ──────────────────────────
def persist_app_state():
    """If LibreOffice is running, send Ctrl+S via xdotool to flush."""
    try:
        import subprocess
        env = os.environ.copy()
        env["DISPLAY"] = ":0"
        subprocess.run(
            ["xdotool", "key", "ctrl+s"],
            env=env, timeout=5,
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
        )
        time.sleep(2)
    except Exception:
        pass

persist_app_state()

# ── locate the pptx file ────────────────────────────────────────────
FILE_PATH = "/home/user/impress_teach_040.pptx"
if not os.path.exists(FILE_PATH):
    # Try common alternate locations
    for p in ["/home/user/Department_Report.pptx",
              "/home/user/Desktop/impress_teach_040.pptx"]:
        if os.path.exists(p):
            FILE_PATH = p
            break

score = 0.0

try:
    from pptx import Presentation
    from pptx.enum.chart import XL_CHART_TYPE

    prs = Presentation(FILE_PATH)

    # ── Check 1: Slide 7 exists and has a chart (0.2) ───────────────
    if len(prs.slides) < 7:
        print(f"Only {len(prs.slides)} slides, need at least 7")
        print(f"REWARD: {score}")
        exit()

    slide = prs.slides[6]  # 0-indexed → slide 7

    chart_shape = None
    for shape in slide.shapes:
        if shape.has_chart:
            chart_shape = shape
            break

    if chart_shape is None:
        print("No chart found on slide 7")
        print(f"REWARD: {score}")
        exit()

    score += 0.2
    print("CHECK 1 PASS: Chart found on slide 7 (+0.2)")

    chart = chart_shape.chart

    # ── Check 2: Chart is a line chart (0.15) ────────────────────────
    # LINE = 4 in XL_CHART_TYPE enum
    is_line = chart.chart_type in (
        XL_CHART_TYPE.LINE,
        XL_CHART_TYPE.LINE_MARKERS,
        XL_CHART_TYPE.LINE_MARKERS_STACKED,
        XL_CHART_TYPE.LINE_MARKERS_STACKED_100,
        XL_CHART_TYPE.LINE_STACKED,
        XL_CHART_TYPE.LINE_STACKED_100,
    )
    if is_line:
        score += 0.15
        print(f"CHECK 2 PASS: Chart is line type ({chart.chart_type}) (+0.15)")
    else:
        print(f"CHECK 2 FAIL: Chart type is {chart.chart_type}, expected LINE")

    # ── Check 3: Correct data values and categories (0.3) ───────────
    expected_categories = ["2020", "2021", "2022", "2023", "2024"]
    expected_values = [150.0, 175.0, 210.0, 195.0, 230.0]

    plot = chart.plots[0]
    actual_cats = list(plot.categories)
    # Categories might be stored as integers in some cases
    actual_cats_str = [str(c).strip() for c in actual_cats]

    cats_match = actual_cats_str == expected_categories

    # Check series values
    vals_match = False
    for series in chart.series:
        actual_vals = list(series.values)
        if len(actual_vals) == len(expected_values):
            if all(abs(a - e) < 0.5 for a, e in zip(actual_vals, expected_values)):
                vals_match = True
                break

    if cats_match and vals_match:
        score += 0.3
        print(f"CHECK 3 PASS: Data values and categories correct (+0.3)")
    elif vals_match:
        score += 0.2
        print(f"CHECK 3 PARTIAL: Values correct but categories wrong (+0.2)")
        print(f"  Expected cats: {expected_categories}, Got: {actual_cats_str}")
    elif cats_match:
        score += 0.1
        print(f"CHECK 3 PARTIAL: Categories correct but values wrong (+0.1)")
    else:
        print(f"CHECK 3 FAIL: Categories={actual_cats_str}, Values don't match")

    # ── Check 4: Chart title is "Enrollment Trends" (0.15) ──────────
    title_text = ""
    try:
        if chart.chart_title and chart.chart_title.has_text_frame:
            title_text = chart.chart_title.text_frame.text.strip()
    except Exception:
        pass

    if title_text == "Enrollment Trends":
        score += 0.15
        print(f"CHECK 4 PASS: Chart title is 'Enrollment Trends' (+0.15)")
    else:
        print(f"CHECK 4 FAIL: Chart title is '{title_text}', expected 'Enrollment Trends'")

    # ── Check 5: Line color is #1565C0 (0.2) ────────────────────────
    color_match = False
    try:
        for series in chart.series:
            line = series.format.line
            if line.color.type is not None:
                rgb = str(line.color.rgb).upper()
                if rgb == "1565C0":
                    color_match = True
                    break
    except Exception as e:
        print(f"  Color check error: {e}")

    if not color_match:
        # Fallback: check XML directly for the color
        try:
            from pptx.oxml.ns import qn
            for series in chart.series:
                ser_xml = series._element.xml
                if "1565C0" in ser_xml.upper() or "1565c0" in ser_xml.lower():
                    color_match = True
                    break
        except Exception:
            pass

    if color_match:
        score += 0.2
        print(f"CHECK 5 PASS: Line color is #1565C0 (+0.2)")
    else:
        print(f"CHECK 5 FAIL: Line color is not #1565C0")

except Exception as e:
    print(f"Error during verification: {e}")
    import traceback
    traceback.print_exc()

# Round to avoid float issues
score = round(score, 2)
print(f"REWARD: {score}")
