"""
initial_setup.py - Create Department_Report.pptx with 9 slides.
Slide 7 has title 'Student Numbers' but NO chart.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

OUTPUT_PATH = "/home/user/impress_teach_040.pptx"


def launch_gui(command: str, delay_sec: float = 1.0):
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            body.paragraphs[0].text = line
        else:
            p = body.add_paragraph()
            p.text = line
            p.level = 0
    return slide


def add_title_only_slide(prs, title_text, extra_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Add title textbox
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    if extra_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, line in enumerate(extra_text):
            if i == 0:
                tf2.paragraphs[0].text = line
            else:
                p2 = tf2.add_paragraph()
                p2.text = line
    return slide


def build_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    add_title_slide(prs, "Department of Computer Science\nAnnual Report 2024",
                    "Prepared by Dr. Sarah Mitchell, Department Chair")

    # Slide 2: Agenda
    add_content_slide(prs, "Agenda", [
        "Faculty Overview and New Hires",
        "Research Output and Publications",
        "Grant Funding Summary",
        "Curriculum Updates",
        "Student Enrollment Numbers",
        "Industry Partnerships",
        "Goals for 2025",
    ])

    # Slide 3: Faculty Overview
    add_content_slide(prs, "Faculty Overview", [
        "Total faculty members: 42 (28 tenured, 14 tenure-track)",
        "New hires in 2024: Dr. James Liu (AI/ML), Dr. Priya Sharma (Cybersecurity)",
        "3 faculty members promoted to Full Professor",
        "Faculty-to-student ratio: 1:18",
        "2 visiting scholars from ETH Zurich and University of Tokyo",
    ])

    # Slide 4: Research Output
    add_content_slide(prs, "Research Output", [
        "127 peer-reviewed publications in top-tier venues",
        "15 papers accepted at NeurIPS, ICML, and CVPR",
        "8 best paper nominations across conferences",
        "3 patents filed for novel algorithms",
        "12 active collaborations with international institutions",
    ])

    # Slide 5: Grant Funding
    add_content_slide(prs, "Grant Funding Summary", [
        "Total grants awarded: $4.2 million",
        "NSF CAREER Awards: 2 new recipients",
        "DOD research contracts: $1.1 million",
        "Industry-sponsored research: $850,000",
        "Internal seed grants distributed: $200,000",
    ])

    # Slide 6: Curriculum Updates
    add_content_slide(prs, "Curriculum Updates", [
        "New course: CS 491 - Large Language Models and Applications",
        "Revised data science track with 3 new electives",
        "Added capstone project requirement for all MS students",
        "Online course offerings expanded to 8 courses",
        "Partnership with AWS Academy for cloud computing certification",
    ])

    # Slide 7: Student Numbers (title only, NO chart)
    add_title_only_slide(prs, "Student Numbers", [
        "This slide presents enrollment data for the department.",
        "Undergraduate and graduate programs have seen consistent growth.",
        "Detailed enrollment figures are available in the supplementary report.",
    ])

    # Slide 8: Industry Partnerships
    add_content_slide(prs, "Industry Partnerships", [
        "Microsoft Research: Joint lab for responsible AI",
        "Google: PhD fellowship program (4 students funded)",
        "Amazon: Cloud infrastructure credits ($50,000)",
        "Intel: Hardware donation for HPC cluster upgrade",
        "Local startups: 15 internship placements for undergraduates",
    ])

    # Slide 9: Goals for 2025
    add_content_slide(prs, "Goals for 2025", [
        "Increase research funding to $5 million",
        "Launch new PhD track in Quantum Computing",
        "Hire 3 additional tenure-track faculty members",
        "Expand industry partnership program",
        "Improve student retention rate to 95%",
    ])

    prs.save(OUTPUT_PATH)
    print(f"Presentation saved to {OUTPUT_PATH}")


if __name__ == "__main__":
    build_presentation()
    launch_gui(f'libreoffice --impress "{OUTPUT_PATH}"', delay_sec=2.0)
    print("LibreOffice Impress launched.")
