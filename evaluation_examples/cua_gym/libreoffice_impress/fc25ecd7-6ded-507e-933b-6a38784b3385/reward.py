"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 25 feels unfinished—how can I wrap it with a 2.00 pt solid border in pure black (#000000) in LibreOffice Impress?
Generated: 2025-09-10 14:08:20
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
import re
from pptx import Presentation
from pptx.util import Emu


def _locate_target_pptx(base_dir: str = "/home/user") -> str | None:
    """Return full path of the *.pptx file that belongs to this task.

    The file name always begins with "slide_25" according to the task
    description, so we search for that pattern in the user directory.
    """
    pattern = re.compile(r"^slide_25.*\.pptx$")
    for fname in os.listdir(base_dir):
        if pattern.match(fname):
            return os.path.join(base_dir, fname)
    return None


def _is_black(rgb_color) -> bool:
    """Return True when pptx RGBColor equals pure black (#000000)."""
    if rgb_color is None:
        return False
    return str(rgb_color).lower() == "000000"


def verify_slide25_border(file_path: str) -> float:
    """Verify that slide 25 has a 2-pt pure-black border around the slide.

    Scoring (progressive):
        • Correct line width   – 0.3
        • Pure black colour    – 0.3
        • Rectangle spans slide – 0.4
    """
    print(f"Checking presentation: {file_path}")

    # Safety checks -----------------------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File not found")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as exc:
        print(f"✗ Could not load PPTX: {exc}")
        print("REWARD: 0.0")
        return 0.0

    if len(prs.slides) < 25:
        print("✗ Presentation contains fewer than 25 slides – cannot verify slide 25")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[24]  # 0-based index
    slide_w, slide_h = prs.slide_width, prs.slide_height
    print(f"Slide dimensions: {slide_w} × {slide_h} EMU")

    # Target metrics -----------------------------------------------------------------
    ONE_POINT_EMU = 12700  # 1 pt in EMU per OOXML spec
    expected_width = 2 * ONE_POINT_EMU  # 2.00 pt → 25 400 EMU
    width_tolerance = 500                # ±0.04 pt tolerance
    position_tolerance = 1000            # 0.1 pt margin around edges

    # Flags for individual requirements
    width_ok = False
    colour_ok = False
    size_ok = False

    # Iterate through shapes on slide 25 ---------------------------------------------
    for shape in slide.shapes:
        if not hasattr(shape, "line") or shape.line is None:
            continue

        line = shape.line
        line_width = line.width or 0

        # 1) Check line width
        if abs(line_width - expected_width) > width_tolerance:
            continue  # not a 2-pt outline – skip further checks on this shape

        width_ok = True  # At least one 2-pt outline exists

        # 2) Check colour is pure black (#000000)
        try:
            if line.color.type == 1 and _is_black(line.color.rgb):
                colour_ok = True
            else:
                continue  # need pure black – otherwise not our border
        except Exception:
            continue  # Unable to read colour – skip shape

        # 3) Check rectangle spans entire slide (within tolerance)
        if (
            shape.left <= position_tolerance and
            shape.top <= position_tolerance and
            shape.width >= slide_w - position_tolerance and
            shape.height >= slide_h - position_tolerance
        ):
            size_ok = True

        # If all three flags are True we found the desired border – can stop early
        if width_ok and colour_ok and size_ok:
            break

    # --------------------------- Scoring --------------------------------------------
    score = 0.0
    if width_ok:
        print("✓ Border width ≈ 2 pt found (0.3)")
        score += 0.3
    else:
        print("✗ No 2 pt border width found (0.0)")

    if colour_ok:
        print("✓ Border colour is pure black (0.3)")
        score += 0.3
    else:
        print("✗ Border colour is not pure black (0.0)")

    if size_ok:
        print("✓ Border rectangle spans the whole slide (0.4)")
        score += 0.4
    else:
        print("✗ Border does not span the whole slide (0.0)")

    # Clamp to [0.0, 1.0]
    score = min(score, 1.0)
    print(f"Total Score: {score}")
    print(f"REWARD: {score}")
    return score


# -----------------------------------------------------------------------------------
# Main execution block – automatically locate the task's PPTX and run verification.
# -----------------------------------------------------------------------------------
if __name__ == "__main__":
    pptx_path = _locate_target_pptx()
    if pptx_path:
        verify_slide25_border(pptx_path)
    else:
        print("✗ Task PPTX file not found in /home/user")
        print("REWARD: 0.0")
