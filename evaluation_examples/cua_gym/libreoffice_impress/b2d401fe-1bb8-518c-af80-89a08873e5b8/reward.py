"""
Reward Script: Move the image on slide 3 to the top area of the slide.
Task ID: osworld_impress_image_top_underline_text_001
Domain: libreoffice_impress
Scoring:
  Component 1 (0.6 pts): Image on slide 3 is in the top area (top < 25% of slide height).
                          Initial: top=2.050in (FAILS), Golden: top=0.300in (PASSES).
  Component 2 (0.4 pts): Image is in the upper-most region (top < 10% of slide height)
                          AND image size is preserved (width/height unchanged).
                          Initial: top=2.050in (FAILS position check), Golden: top=0.300in (PASSES).
Total: 1.0
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_image_top_underline_text_001'

# Known initial image dimensions from setup (EMU units)
# These are expected to be preserved in the golden file
EXPECTED_IMAGE_WIDTH  = 4114800   # 4.5 inches
EXPECTED_IMAGE_HEIGHT = 3108960   # 3.4 inches

# Slide 3 is index 2 (0-based)
TARGET_SLIDE_INDEX = 2


def verify_task(file_path):
    """
    Verify that the image on slide 3 has been moved to the top area of the slide.
    Returns a float between 0.0 and 1.0.
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: verify 5-slide deck exists and slide 3 is accessible
    try:
        if len(prs.slides) < 3:
            print(f"FAIL: Expected at least 3 slides, found {len(prs.slides)}")
            print("REWARD: 0.0")
            return 0.0
    except Exception as e:
        print(f"CRITICAL: Cannot check slide count: {e}")
        print("REWARD: 0.0")
        return 0.0

    slide_height = prs.slide_height  # 6858000 EMU = 7.5 inches
    slide = prs.slides[TARGET_SLIDE_INDEX]

    # Locate the image shape on slide 3
    image_shape = None
    try:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE = 13
                image_shape = shape
                break
    except Exception as e:
        print(f"ERROR: Could not iterate slide 3 shapes: {e}")

    if image_shape is None:
        print("FAIL: No image (PICTURE shape) found on slide 3")
        print(f"REWARD: 0.0")
        return 0.0

    img_top    = image_shape.top
    img_width  = image_shape.width
    img_height = image_shape.height

    print(f"INFO: Image found on slide 3: top={img_top} ({img_top/914400:.3f}in), "
          f"width={img_width}, height={img_height}")
    print(f"INFO: Slide height={slide_height} EMU ({slide_height/914400:.2f}in)")

    # Component 1: Image is in the top area of the slide (top < 25% of slide height)
    # top_25pct = 1,714,500 EMU = 1.875 inches
    # Initial top = 1,874,520 EMU (2.050 in) → FAILS (above 25% threshold)
    # Golden top  =   274,320 EMU (0.300 in) → PASSES
    try:
        top_25pct_threshold = slide_height * 0.25  # 1,714,500 EMU
        if img_top < top_25pct_threshold:
            print(f"PASS: Component 1 — Image is in the top area: "
                  f"top={img_top} EMU ({img_top/914400:.3f}in) < threshold={top_25pct_threshold:.0f} EMU "
                  f"({top_25pct_threshold/914400:.3f}in) (0.6 pts)")
            total_score += 0.6
        else:
            print(f"FAIL: Component 1 — Image NOT in top area: "
                  f"top={img_top} EMU ({img_top/914400:.3f}in) >= threshold={top_25pct_threshold:.0f} EMU "
                  f"({top_25pct_threshold/914400:.3f}in)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Image is in the upper-most region (top < 10% of slide height)
    # AND image dimensions are preserved (width and height unchanged)
    # top_10pct = 685,800 EMU = 0.75 inches
    # Initial top = 1,874,520 EMU → FAILS position check
    # Golden top  =   274,320 EMU → PASSES (0.300in < 0.75in)
    try:
        top_10pct_threshold = slide_height * 0.10  # 685,800 EMU
        position_ok = img_top < top_10pct_threshold
        width_ok  = img_width  == EXPECTED_IMAGE_WIDTH
        height_ok = img_height == EXPECTED_IMAGE_HEIGHT

        if position_ok and width_ok and height_ok:
            print(f"PASS: Component 2 — Image is in upper-most region "
                  f"(top={img_top/914400:.3f}in < {top_10pct_threshold/914400:.3f}in) "
                  f"AND dimensions preserved ({img_width}x{img_height}) (0.4 pts)")
            total_score += 0.4
        else:
            msgs = []
            if not position_ok:
                msgs.append(f"top={img_top/914400:.3f}in >= threshold={top_10pct_threshold/914400:.3f}in")
            if not width_ok:
                msgs.append(f"width changed: got {img_width}, expected {EXPECTED_IMAGE_WIDTH}")
            if not height_ok:
                msgs.append(f"height changed: got {img_height}, expected {EXPECTED_IMAGE_HEIGHT}")
            print(f"FAIL: Component 2 — {'; '.join(msgs)}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
