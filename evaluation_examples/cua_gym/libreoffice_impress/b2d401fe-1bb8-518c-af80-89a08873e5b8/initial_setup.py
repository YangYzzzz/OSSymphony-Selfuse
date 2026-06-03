"""
Initial Setup: 5-slide product catalog deck; slide 3 has an image in the center.
Task ID: osworld_impress_image_top_underline_text_001
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image as PILImage, ImageDraw, ImageFont

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_image_top_underline_text_001'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
IMG_PATH = f'{WORKDIR}/{TASK_ID}_product.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_product_image(path: str, label: str, bg_color: tuple, text_color: tuple = (255, 255, 255)):
    """Create a simple product image with a label."""
    width, height = 400, 300
    img = PILImage.new('RGB', (width, height), bg_color)
    draw = ImageDraw.Draw(img)
    # Draw border
    draw.rectangle([4, 4, width - 5, height - 5], outline=text_color, width=3)
    # Draw text in center
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf", 36)
        small_font = ImageFont.truetype("/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf", 20)
    except Exception:
        font = ImageFont.load_default()
        small_font = font
    # Center the main label
    bbox = draw.textbbox((0, 0), label, font=font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    draw.text(((width - tw) // 2, (height - th) // 2 - 20), label, font=font, fill=text_color)
    # Subtitle
    sub = "Product Photo"
    bbox2 = draw.textbbox((0, 0), sub, font=small_font)
    sw = bbox2[2] - bbox2[0]
    draw.text(((width - sw) // 2, (height - th) // 2 + 30), sub, font=small_font, fill=text_color)
    img.save(path)


def add_text_box(slide, text, left, top, width, height, font_size=18, bold=False,
                 color=(0x33, 0x33, 0x33), alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return txBox


def create_initial():
    prs = Presentation()
    # Standard widescreen 10x7.5 inches
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_layout_blank = prs.slide_layouts[5]   # Blank
    slide_layout_title = prs.slide_layouts[0]   # Title Slide

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(slide_layout_title)
    slide1.shapes.title.text = "TechVision Pro Product Catalog 2025"
    slide1.placeholders[1].text = "Innovative Solutions for Modern Enterprises"
    # Style title
    for run in slide1.shapes.title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
    bg1 = slide1.background.fill
    bg1.solid()
    bg1.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)

    # ---- Slide 2: Product Overview ----
    slide2 = prs.slides.add_slide(slide_layout_blank)
    bg2 = slide2.background.fill
    bg2.solid()
    bg2.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    add_text_box(slide2, "Product Overview", Inches(0.5), Inches(0.3),
                 Inches(9), Inches(0.8), font_size=28, bold=True,
                 color=(0x1F, 0x49, 0x7D), alignment=PP_ALIGN.LEFT)

    overview_text = (
        "TechVision Pro offers a comprehensive range of enterprise-grade solutions tailored "
        "to meet the evolving needs of modern businesses. Our product lineup includes hardware, "
        "software, and hybrid solutions with best-in-class support.\n\n"
        "• CloudSync X1 — Enterprise Cloud Storage\n"
        "• DataShield Pro — Advanced Cybersecurity Suite\n"
        "• VisionBoard 4K — Interactive Display System\n"
        "• ConnectHub Ultra — Next-Gen Networking\n"
        "• AutoDesk AI — Intelligent Workflow Automation"
    )
    add_text_box(slide2, overview_text, Inches(0.5), Inches(1.3),
                 Inches(9), Inches(5.5), font_size=16,
                 color=(0x33, 0x33, 0x33), alignment=PP_ALIGN.LEFT)

    # ---- Slide 3: VisionBoard 4K — image in CENTER ----
    slide3 = prs.slides.add_slide(slide_layout_blank)
    bg3 = slide3.background.fill
    bg3.solid()
    bg3.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Title at top
    add_text_box(slide3, "VisionBoard 4K — Interactive Display System",
                 Inches(0.5), Inches(0.2), Inches(9), Inches(0.7),
                 font_size=24, bold=True, color=(0x1F, 0x49, 0x7D))

    # Product image positioned in the CENTER of the slide
    create_product_image(IMG_PATH, "VisionBoard 4K", (0x2E, 0x86, 0xAB))
    img_width = Inches(4.5)
    img_height = Inches(3.4)
    img_left = (prs.slide_width - img_width) // 2      # horizontally centered
    img_top = (prs.slide_height - img_height) // 2     # vertically centered (center of slide)
    slide3.shapes.add_picture(IMG_PATH, img_left, img_top, img_width, img_height)

    # Caption below image
    add_text_box(slide3, "Resolution: 3840x2160 | Touch: 20-point | Connectivity: HDMI 2.1, USB-C, Wi-Fi 6",
                 Inches(0.5), Inches(6.5), Inches(9), Inches(0.7),
                 font_size=13, color=(0x55, 0x55, 0x55), alignment=PP_ALIGN.CENTER)

    # ---- Slide 4: DataShield Pro ----
    slide4 = prs.slides.add_slide(slide_layout_blank)
    bg4 = slide4.background.fill
    bg4.solid()
    bg4.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    add_text_box(slide4, "DataShield Pro — Advanced Cybersecurity Suite",
                 Inches(0.5), Inches(0.2), Inches(9), Inches(0.7),
                 font_size=24, bold=True, color=(0x1F, 0x49, 0x7D))

    features = (
        "Key Features:\n"
        "• Real-time threat detection powered by AI\n"
        "• Zero-trust architecture with identity verification\n"
        "• End-to-end encryption for data at rest and in transit\n"
        "• Compliance reporting: GDPR, HIPAA, SOC 2\n"
        "• Automated incident response with 99.9% uptime SLA\n\n"
        "Pricing: Starting at $2,400/year per 50 endpoints\n"
        "Support: 24/7 dedicated security operations center"
    )
    add_text_box(slide4, features, Inches(0.5), Inches(1.2),
                 Inches(9), Inches(5.5), font_size=16,
                 color=(0x33, 0x33, 0x33))

    # ---- Slide 5: Contact & Next Steps ----
    slide5 = prs.slides.add_slide(slide_layout_blank)
    bg5 = slide5.background.fill
    bg5.solid()
    bg5.fore_color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    add_text_box(slide5, "Get Started Today", Inches(1), Inches(1.5),
                 Inches(8), Inches(1.0), font_size=36, bold=True,
                 color=(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)

    contact_info = (
        "Contact Our Sales Team:\n\n"
        "Email: enterprise@techvisionpro.com\n"
        "Phone: +1 (800) 555-0192\n"
        "Website: www.techvisionpro.com\n\n"
        "Schedule a free 30-minute demo with one of our specialists."
    )
    add_text_box(slide5, contact_info, Inches(1.5), Inches(2.8),
                 Inches(7), Inches(4.0), font_size=18,
                 color=(0xE8, 0xF0, 0xFE), alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
