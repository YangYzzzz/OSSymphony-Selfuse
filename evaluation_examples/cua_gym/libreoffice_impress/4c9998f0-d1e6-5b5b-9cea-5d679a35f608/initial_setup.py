"""
Initial Setup: Sociology presentation with 8 slides, default line spacing
Task ID: impress_stu_015
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_015'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_bulleted_text(text_frame, items, font_size=Pt(18), font_name="Arial"):
    """Add bulleted text items to a text frame with default single line spacing."""
    text_frame.clear()
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.size = font_size
            run.font.name = font_name
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def add_sub_bullets(text_frame, items, font_size=Pt(16), font_name="Arial"):
    """Add sub-bulleted text items at indent level 1."""
    for item in items:
        p = text_frame.add_paragraph()
        p.text = item
        p.level = 1
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.size = font_size
            run.font.name = font_name
            run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Introduction to Sociology"
    slide1.placeholders[1].text = "Understanding Social Structures and Human Behavior\nProfessor Elena Vasquez | SOC 201 | Spring 2026"

    # --- Slide 2: What is Sociology? ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "What is Sociology?"
    tf2 = slide2.placeholders[1].text_frame
    add_bulleted_text(tf2, [
        "The systematic study of human society and social interaction",
        "Examines how groups, institutions, and structures shape behavior",
        "Uses empirical research methods to analyze social phenomena",
        "Founded as a discipline in the 19th century by Auguste Comte",
        "Bridges the gap between individual experience and broader social forces",
    ])

    # --- Slide 3: Key Sociological Perspectives ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Key Sociological Perspectives"
    tf3 = slide3.placeholders[1].text_frame
    add_bulleted_text(tf3, [
        "Structural Functionalism: Society as an interconnected system",
        "Conflict Theory: Power struggles and inequality drive social change",
        "Symbolic Interactionism: Meaning is constructed through daily interaction",
        "Feminist Theory: Gender as a central axis of social organization",
    ])
    add_sub_bullets(tf3, [
        "Each perspective offers a unique lens for analyzing social issues",
        "Modern sociology often combines multiple perspectives",
    ])

    # --- Slide 4: Social Stratification ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Social Stratification"
    tf4 = slide4.placeholders[1].text_frame
    add_bulleted_text(tf4, [
        "Hierarchical arrangement of individuals into social categories",
        "Class, race, gender, and age are primary stratification dimensions",
        "Weber identified class, status, and party as key dimensions of power",
        "Social mobility varies significantly across different societies",
        "Intersectionality reveals how overlapping identities shape experiences",
    ])

    # --- Slide 5: Research Methods in Sociology ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Research Methods in Sociology"
    tf5 = slide5.placeholders[1].text_frame
    add_bulleted_text(tf5, [
        "Surveys and questionnaires for large-scale data collection",
        "Ethnography and participant observation for in-depth understanding",
        "Content analysis of media, documents, and cultural artifacts",
        "Statistical analysis to identify patterns and correlations",
        "Mixed methods approaches combine quantitative and qualitative data",
    ])

    # --- Slide 6: Socialization and Identity ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Socialization and Identity"
    tf6 = slide6.placeholders[1].text_frame
    add_bulleted_text(tf6, [
        "Socialization is the lifelong process of learning cultural norms",
        "Primary agents: family, peers, education, and media",
        "Erving Goffman's dramaturgical approach to self-presentation",
        "Identity formation involves negotiation between self and society",
    ])

    # --- Slide 7: Deviance and Social Control ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Deviance and Social Control"
    tf7 = slide7.placeholders[1].text_frame
    add_bulleted_text(tf7, [
        "Deviance is behavior that violates established social norms",
        "Labeling theory: deviance is socially constructed through reactions",
        "Formal controls include laws, regulations, and institutional rules",
        "Informal controls include peer pressure, ridicule, and ostracism",
        "Durkheim argued some deviance is functional for society",
    ])

    # --- Slide 8: Discussion Questions ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Discussion Questions"
    tf8 = slide8.placeholders[1].text_frame
    add_bulleted_text(tf8, [
        "How do social media platforms reshape socialization processes?",
        "In what ways does your own social position influence your worldview?",
        "Can sociology be truly objective, or is it inherently political?",
        "What role should sociological research play in public policy?",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
