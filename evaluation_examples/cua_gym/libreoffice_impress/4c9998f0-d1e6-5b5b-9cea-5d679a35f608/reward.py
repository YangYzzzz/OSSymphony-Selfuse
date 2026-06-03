"""
Reward Script: Set line spacing of bullet point text on slides 2-5 to 1.5 lines
Task ID: impress_stu_015
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Slide 2 bullet text has 1.5 line spacing
  Component 2 (0.25): Slide 3 bullet text has 1.5 line spacing
  Component 3 (0.25): Slide 4 bullet text has 1.5 line spacing
  Component 4 (0.25): Slide 5 bullet text has 1.5 line spacing
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_015'


def get_bullet_paragraphs(slide):
    """Get non-title text paragraphs from a slide (bullet point text).

    Title shapes typically have placeholder index 0 (title) or are the first
    shape with text. We identify bullet content shapes as non-title shapes
    with text frames containing multiple paragraphs of content.
    """
    bullet_paras = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # Skip title placeholders
        if shape.is_placeholder:
            ph_idx = shape.placeholder_format.idx
            # idx 0 = title, idx 1 = subtitle (on title slides)
            # For content slides, idx 0 = title, idx 1 = content body
            if ph_idx == 0:
                continue
        # Collect non-empty paragraphs from content shapes
        for para in shape.text_frame.paragraphs:
            if para.text.strip():
                bullet_paras.append(para)
    return bullet_paras


def check_slide_line_spacing(slide, slide_num, expected_spacing=1.5):
    """Check if bullet paragraphs on a slide have the expected line spacing.

    Returns (score_fraction, total_checked, passed_count)
    """
    bullet_paras = get_bullet_paragraphs(slide)
    if not bullet_paras:
        print(f"  WARNING: Slide {slide_num} has no bullet paragraphs")
        return 0.0, 0, 0

    passed = 0
    total = len(bullet_paras)
    for para in bullet_paras:
        actual = para.line_spacing
        if actual is not None and abs(actual - expected_spacing) < 0.01:
            passed += 1
        else:
            print(f"  DETAIL: Slide {slide_num} para '{para.text[:40]}...' "
                  f"spacing={actual} (expected {expected_spacing})")

    if total > 0:
        fraction = passed / total
    else:
        fraction = 0.0
    return fraction, total, passed


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Verify we have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Expected at least 5 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Components 1-4: Check slides 2-5 (indices 1-4) for 1.5 line spacing
    target_slides = [2, 3, 4, 5]
    weight_per_slide = 0.25

    for slide_num in target_slides:
        slide_idx = slide_num - 1
        component_num = slide_num - 1  # Components 1-4
        try:
            slide = prs.slides[slide_idx]
            fraction, total, passed = check_slide_line_spacing(slide, slide_num, 1.5)

            if fraction >= 1.0:
                print(f"PASS: Component {component_num} -- Slide {slide_num}: "
                      f"all {passed}/{total} bullet paragraphs have 1.5 line spacing "
                      f"({weight_per_slide} pts)")
                total_score += weight_per_slide
            elif fraction > 0.0:
                partial = weight_per_slide * fraction
                print(f"PARTIAL: Component {component_num} -- Slide {slide_num}: "
                      f"{passed}/{total} bullet paragraphs have 1.5 line spacing "
                      f"({partial:.3f} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component {component_num} -- Slide {slide_num}: "
                      f"0/{total} bullet paragraphs have 1.5 line spacing")
        except Exception as e:
            print(f"ERROR: Component {component_num} -- Slide {slide_num}: {e}")

    final_score = min(round(total_score, 4), 1.0)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
