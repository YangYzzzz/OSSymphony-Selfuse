"""
FINAL REWARD SCRIPT - SUCCESS
Task: Create a user field 'ProjectCode' = 'WR-2025-09' and put it in the header.
Generated: 2025-10-17 07:45:57
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation

FILE_PATH = "/home/user/create_a_user_field_projectcode_wr_2025_09_and_put_it_in_the_header.pptx"

# -----------------------------------------------------------------------------
# Helper: Verify a custom document property exists and equals the expected value
# -----------------------------------------------------------------------------

def check_custom_property(file_path: str, property_name: str, expected_value: str) -> bool:
    """Return True iff the PPTX contains a custom property with the given value."""
    try:
        with zipfile.ZipFile(file_path, 'r') as z:
            # custom properties live in this fixed location inside a PPTX package
            custom_path = 'docProps/custom.xml'
            if custom_path not in z.namelist():
                print("✗ custom.xml not found – no custom properties present")
                return False
            xml_content = z.read(custom_path)

        # parse XML safely – namespaces are required for correct element access
        ns = {
            'cp': 'http://schemas.openxmlformats.org/officeDocument/2006/custom-properties',
            'vt': 'http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes'
        }
        root = ET.fromstring(xml_content)

        # iterate through <property> elements looking for the right name
        for prop in root.findall('cp:property', ns):
            if prop.get('name') == property_name:
                # first (and only) child holds the value (e.g. <vt:lpwstr>)
                value_element = list(prop)[0]
                actual_value = value_element.text
                print(f"Found custom property '{property_name}' with value '{actual_value}'")
                if actual_value == expected_value:
                    print("✓ Custom property value matches expected")
                    return True
                else:
                    print("✗ Custom property value does NOT match expected")
                    return False

        # property name not found at all
        print(f"✗ Custom property '{property_name}' not found in file")
        return False

    except Exception as e:
        print(f"✗ Error while reading custom properties: {e}")
        return False

# -----------------------------------------------------------------------------
# Helper: Verify the expected header text appears on slides
# -----------------------------------------------------------------------------

def check_header_text(file_path: str, expected_text: str):
    """Return (any_slide_has_text, all_slides_have_text)."""
    try:
        prs = Presentation(file_path)

        any_slide_has_text = False
        all_slides_have_text = True  # will be flipped to False if any slide misses it

        for idx, slide in enumerate(prs.slides, start=1):
            slide_contains = False
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text and expected_text in shape.text:
                    # Found the header text on this slide
                    slide_contains = True
                    # Log vertical position if available – helpful for debug
                    try:
                        top_inches = shape.top / 914400  # EMUs → inches
                        print(f"✓ Found expected text on slide {idx} (top={top_inches:.2f} in)")
                    except Exception:
                        print(f"✓ Found expected text on slide {idx}")
                    break  # no need to inspect other shapes on this slide

            if slide_contains:
                any_slide_has_text = True
            else:
                all_slides_have_text = False
                print(f"✗ Expected text NOT found on slide {idx}")

        return any_slide_has_text, all_slides_have_text

    except Exception as e:
        print(f"✗ Error while checking slide headers: {e}")
        return False, False

# -----------------------------------------------------------------------------
# Main verification entry point
# -----------------------------------------------------------------------------

def verify_task(file_path: str) -> float:
    """Return a progressive score (0.0 – 1.0) based on task completion."""
    print(f"Verifying PPTX task for file: {file_path}")

    max_score = 1.0
    score = 0.0

    # Requirement 1: Correct custom property (60%)
    if check_custom_property(file_path, 'ProjectCode', 'WR-2025-09'):
        score += 0.6
    else:
        print("Custom property requirement failed – 0 points")

    # Requirement 2: Header text present (40%) – partial & full credit
    any_slide, all_slides = check_header_text(file_path, 'WR-2025-09')

    # 0.2 if it appears on at least one slide, +0.2 bonus if EVERY slide has it
    if any_slide:
        score += 0.2
    if any_slide and all_slides:
        score += 0.2

    final_score = min(score, max_score)
    print(f"Total Score Awarded: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score

# -----------------------------------------------------------------------------
# Execute verification when the script is run directly
# -----------------------------------------------------------------------------
if __name__ == "__main__":
    verify_task(FILE_PATH)

