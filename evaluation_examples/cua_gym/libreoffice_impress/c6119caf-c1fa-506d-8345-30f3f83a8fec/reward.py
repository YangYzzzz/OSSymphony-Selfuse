"""
Reward Script: Rename slide 3 title to 'Key Findings' and add speaker note
Task ID: osworld_impress_multi_op_combined_003
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 3 title is 'Key Findings'        — 0.6 points
  Component 2: Slide 3 speaker note is correct         — 0.4 points
  Total: 1.0
"""

import os

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_multi_op_combined_003'

EXPECTED_TITLE = 'Key Findings'
EXPECTED_NOTE = 'Emphasize data quality here'
SLIDE_INDEX = 2  # 0-based index for slide 3


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.

    Task requires:
    1. The title of slide 3 changed from 'Results' to 'Key Findings'
    2. Slide 3 has the speaker note 'Emphasize data quality here' added

    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation file
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"CRITICAL: Presentation has fewer than 3 slides ({len(prs.slides)} slides found)")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[SLIDE_INDEX]

    # Component 1: Slide 3 title changed to 'Key Findings' (0.6 points)
    # This FAILS on initial (title='Results') and PASSES on golden (title='Key Findings')
    try:
        title_shape = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name and 'Title' in shape.name:
                title_shape = shape
                break

        if title_shape is None:
            print("FAIL: Component 1 — No title shape found on slide 3")
        else:
            actual_title = title_shape.text_frame.paragraphs[0].text.strip()
            if actual_title == EXPECTED_TITLE:
                print(f"PASS: Component 1 — Slide 3 title is '{actual_title}' (0.6 pts)")
                total_score += 0.6
            else:
                print(f"FAIL: Component 1 — Expected title '{EXPECTED_TITLE}', found '{actual_title}'")
    except Exception as e:
        print(f"ERROR: Component 1 — Could not check slide 3 title: {e}")

    # Component 2: Slide 3 speaker note is 'Emphasize data quality here' (0.4 points)
    # This FAILS on initial (no notes) and PASSES on golden (note added)
    try:
        notes_text = ""
        try:
            notes_tf = slide.notes_slide.notes_text_frame
            # notes_text_frame includes a copy of the slide text in some cases;
            # iterate paragraphs and collect non-title text
            paras = [p.text.strip() for p in notes_tf.paragraphs if p.text.strip()]
            # The notes body placeholder text is what the user entered;
            # the first placeholder in notes_slide is the slide image reference,
            # the second is the notes body. We gather all text and check for expected note.
            notes_text = " ".join(paras)
        except Exception:
            notes_text = ""

        if EXPECTED_NOTE in notes_text:
            print(f"PASS: Component 2 — Speaker note contains '{EXPECTED_NOTE}' (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 2 — Expected note '{EXPECTED_NOTE}', found notes text: '{notes_text}'")
    except Exception as e:
        print(f"ERROR: Component 2 — Could not check slide 3 speaker notes: {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in golden/initial env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
