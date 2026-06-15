"""
Initial Setup: 6-slide research summary presentation
Task ID: osworld_impress_multi_op_combined_003
Domain: libreoffice_impress

Slide 3 has title 'Results' and NO speaker notes.
Task: Rename slide 3 title to 'Key Findings' and add speaker note.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_multi_op_combined_003'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Default slide size: 10 inches wide, 7.5 inches tall
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_layouts = prs.slide_layouts
    # Layout 0: Title Slide, Layout 1: Title+Content, Layout 5: Blank, Layout 6: Title Only

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(slide_layouts[0])
    slide1.shapes.title.text = "Climate Change Impact on Coastal Ecosystems"
    slide1.placeholders[1].text = "A Comprehensive Research Summary\nDr. Elena Marchetti et al.\nOcean Sciences Institute, 2024"

    # --- Slide 2: Introduction ---
    slide2 = prs.slides.add_slide(slide_layouts[1])
    slide2.shapes.title.text = "Introduction"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Research Background"
    p2 = tf2.add_paragraph()
    p2.text = "Coastal ecosystems face unprecedented pressure from rising sea levels and temperature anomalies."
    p2.level = 1
    p3 = tf2.add_paragraph()
    p3.text = "This study examines 15 years of longitudinal data across 12 monitoring stations."
    p3.level = 1
    p4 = tf2.add_paragraph()
    p4.text = "Data collected between 2008 and 2023 from the Pacific Northwest coastline."
    p4.level = 1
    slide2.notes_slide.notes_text_frame.text = "Introduce study objectives and team members."

    # --- Slide 3: Results (INITIAL — title must be 'Results', NO speaker notes) ---
    slide3 = prs.slides.add_slide(slide_layouts[1])
    slide3.shapes.title.text = "Results"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Core Findings"
    p3a = tf3.add_paragraph()
    p3a.text = "Sea surface temperature increased by 1.4°C over the study period."
    p3a.level = 1
    p3b = tf3.add_paragraph()
    p3b.text = "Kelp forest coverage declined by 38% in northern monitoring zones."
    p3b.level = 1
    p3c = tf3.add_paragraph()
    p3c.text = "Salinity variance exceeded historical baselines by a factor of 2.3."
    p3c.level = 1
    p3d = tf3.add_paragraph()
    p3d.text = "Species diversity index dropped from 4.7 to 3.1 (Shannon-Wiener)."
    p3d.level = 1
    # NOTE: NO speaker notes added to slide 3 — this is the initial state

    # --- Slide 4: Methodology ---
    slide4 = prs.slides.add_slide(slide_layouts[1])
    slide4.shapes.title.text = "Methodology"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Data Collection and Analysis"
    p4a = tf4.add_paragraph()
    p4a.text = "Continuous underwater sensor arrays at 12 coastal stations."
    p4a.level = 1
    p4b = tf4.add_paragraph()
    p4b.text = "Monthly biological surveys using transect and quadrat methods."
    p4b.level = 1
    p4c = tf4.add_paragraph()
    p4c.text = "Statistical modeling via generalized additive mixed models (GAMM)."
    p4c.level = 1
    slide4.notes_slide.notes_text_frame.text = "Highlight the robustness of the sensor network."

    # --- Slide 5: Discussion ---
    slide5 = prs.slides.add_slide(slide_layouts[1])
    slide5.shapes.title.text = "Discussion"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Interpretation of Results"
    p5a = tf5.add_paragraph()
    p5a.text = "Temperature anomalies correlate strongly (r=0.87) with biodiversity loss."
    p5a.level = 1
    p5b = tf5.add_paragraph()
    p5b.text = "Kelp decline accelerates positive feedback loops in ocean acidification."
    p5b.level = 1
    p5c = tf5.add_paragraph()
    p5c.text = "Projected losses of 60-75% of endemic species by 2050 if trends continue."
    p5c.level = 1

    # --- Slide 6: Conclusion ---
    slide6 = prs.slides.add_slide(slide_layouts[1])
    slide6.shapes.title.text = "Conclusion"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Summary and Recommendations"
    p6a = tf6.add_paragraph()
    p6a.text = "Immediate policy intervention required to protect critical habitats."
    p6a.level = 1
    p6b = tf6.add_paragraph()
    p6b.text = "Expanded monitoring network recommended for real-time risk assessment."
    p6b.level = 1
    p6c = tf6.add_paragraph()
    p6c.text = "Collaborative conservation strategy with federal and local agencies is essential."
    p6c.level = 1
    slide6.notes_slide.notes_text_frame.text = "Thank the funding agencies and field teams."

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
