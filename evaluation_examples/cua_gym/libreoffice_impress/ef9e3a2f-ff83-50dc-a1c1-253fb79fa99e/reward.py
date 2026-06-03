"""
Reward Script: Nonprofit fundraising presentation with 10 slides
Task ID: impress_wf_063
Domain: libreoffice_impress
Scoring:
  C1: File on Desktop + 10 slides (0.15)
  C2: Slide 1 title text (0.10)
  C3: Slide 2 heart shape with grow animation (0.15)
  C4: Slide 3 has 3 counter number elements (0.10)
  C5: Slide 4 has 4 program card shapes (0.10)
  C6: Slide 6 has a pie chart (0.15)
  C7: Slide 7 has donation tiers table (0.10)
  C8: Slide 10 has QR placeholder square shape (0.05)
  C9: Orange #E65100 accent color used across slides (0.10)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_063'
FILE_PATH = os.path.join(WORKDIR, 'Desktop', 'Fundraiser.pptx')


def get_all_text(slide):
    """Recursively get all text from shapes including groups."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                texts.append(para.text.strip())
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                if hasattr(sub, 'text_frame') and sub.has_text_frame:
                    for para in sub.text_frame.paragraphs:
                        texts.append(para.text.strip())
    return texts


def check_animation_on_slide(pptx_path, slide_num):
    """Check if a slide has animation elements (1-indexed slide_num)."""
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            fname = f'ppt/slides/slide{slide_num}.xml'
            with zf.open(fname) as f:
                content = f.read().decode()
                # animScale is the element for Grow/Shrink animation
                return 'animScale' in content or 'anim:' in content or '<p:anim' in content
    except Exception:
        return False


def get_font_colors_from_slide(slide):
    """Extract all font RGB colors from a slide's text runs."""
    colors = set()
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.color.type is not None:
                            colors.add(str(run.font.color.rgb).upper())
                    except Exception:
                        pass
    return colors


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: File exists on Desktop and has 10 slides (0.15 pts)
    try:
        if num_slides == 10:
            print(f"PASS: Component 1 — File has exactly 10 slides (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — Expected 10 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    if num_slides < 10:
        # Not enough slides for the rest of the checks
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: Slide 1 has "Hope Foundation" title text (0.10 pts)
    try:
        slide1_texts = get_all_text(prs.slides[0])
        combined = " ".join(slide1_texts).lower()
        has_hope = "hope foundation" in combined
        has_fundraiser = "fundraiser" in combined or "annual" in combined
        if has_hope and has_fundraiser:
            print(f"PASS: Component 2 — Slide 1 has 'Hope Foundation' and fundraiser text (0.10 pts)")
            total_score += 0.10
        elif has_hope:
            print(f"PARTIAL: Component 2 — Slide 1 has 'Hope Foundation' but missing fundraiser text (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 2 — Slide 1 missing 'Hope Foundation'. Found: {slide1_texts[:3]}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 2 has a heart shape with Grow animation (0.15 pts)
    try:
        slide2 = prs.slides[1]
        heart_found = False
        for shape in slide2.shapes:
            name_lower = shape.name.lower()
            if 'heart' in name_lower:
                heart_found = True
                break
            # Also check if it's an auto shape that might be a heart
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # Check via XML for heart preset
                try:
                    xml_str = shape._element.xml
                    if 'heart' in xml_str.lower():
                        heart_found = True
                        break
                except Exception:
                    pass

        has_anim = check_animation_on_slide(file_path, 2)

        if heart_found and has_anim:
            print(f"PASS: Component 3 — Slide 2 has heart shape with animation (0.15 pts)")
            total_score += 0.15
        elif heart_found:
            print(f"PARTIAL: Component 3 — Slide 2 has heart but no animation (0.08 pts)")
            total_score += 0.08
        elif has_anim:
            print(f"PARTIAL: Component 3 — Slide 2 has animation but no heart shape found (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 3 — Slide 2 missing heart shape and animation")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 3 has 3 counter/number text elements (0.10 pts)
    try:
        slide3 = prs.slides[2]
        # Count shapes that contain large numbers (counters)
        counter_shapes = 0
        for shape in slide3.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                # Counter shapes contain numbers like "12,450" or "285,000"
                cleaned = text.replace(',', '').replace('.', '')
                if cleaned.isdigit() and len(cleaned) >= 3:
                    counter_shapes += 1

        if counter_shapes >= 3:
            print(f"PASS: Component 4 — Slide 3 has {counter_shapes} counter elements (0.10 pts)")
            total_score += 0.10
        elif counter_shapes >= 2:
            print(f"PARTIAL: Component 4 — Slide 3 has {counter_shapes}/3 counters (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 4 — Slide 3 has {counter_shapes} counter elements, expected 3")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slide 4 has 4 program card shapes (rounded rectangles or similar) (0.10 pts)
    try:
        slide4 = prs.slides[3]
        card_shapes = 0
        for shape in slide4.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # Check if it has text content (program cards have text)
                if shape.has_text_frame and shape.text_frame.text.strip():
                    card_shapes += 1

        if card_shapes >= 4:
            print(f"PASS: Component 5 — Slide 4 has {card_shapes} program card shapes (0.10 pts)")
            total_score += 0.10
        elif card_shapes >= 2:
            print(f"PARTIAL: Component 5 — Slide 4 has {card_shapes}/4 cards (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 5 — Slide 4 has {card_shapes} card shapes, expected 4")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Slide 6 has a pie chart (0.15 pts)
    try:
        slide6 = prs.slides[5]
        chart_found = False
        pie_chart = False
        for shape in slide6.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART or shape.shape_type == 3:
                chart_found = True
                try:
                    chart = shape.chart
                    # PIE chart type is 5 in python-pptx
                    if chart.chart_type == 5 or 'PIE' in str(chart.chart_type).upper():
                        pie_chart = True
                except Exception:
                    pass

        if pie_chart:
            print(f"PASS: Component 6 — Slide 6 has a pie chart (0.15 pts)")
            total_score += 0.15
        elif chart_found:
            print(f"PARTIAL: Component 6 — Slide 6 has a chart but not pie type (0.08 pts)")
            total_score += 0.08
        else:
            print(f"FAIL: Component 6 — Slide 6 has no chart")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Slide 7 has a donation tiers table with Bronze/Silver/Gold/Platinum (0.10 pts)
    try:
        slide7 = prs.slides[6]
        table_found = False
        tiers_found = set()
        expected_tiers = {'bronze', 'silver', 'gold', 'platinum'}

        for shape in slide7.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE or shape.shape_type == 19:
                table_found = True
                table = shape.table
                for row_idx in range(len(table.rows)):
                    for col_idx in range(len(table.columns)):
                        cell_text = table.cell(row_idx, col_idx).text.strip().lower()
                        for tier in expected_tiers:
                            if tier in cell_text:
                                tiers_found.add(tier)

        if table_found and len(tiers_found) >= 4:
            print(f"PASS: Component 7 — Slide 7 has table with all tiers: {tiers_found} (0.10 pts)")
            total_score += 0.10
        elif table_found and len(tiers_found) >= 2:
            print(f"PARTIAL: Component 7 — Slide 7 has table with {len(tiers_found)}/4 tiers (0.05 pts)")
            total_score += 0.05
        elif table_found:
            print(f"PARTIAL: Component 7 — Slide 7 has table but missing tier names (0.03 pts)")
            total_score += 0.03
        else:
            print(f"FAIL: Component 7 — Slide 7 has no table")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    # Component 8: Slide 10 has a QR placeholder square shape (0.05 pts)
    try:
        slide10 = prs.slides[9]
        qr_found = False
        for shape in slide10.shapes:
            # Check for square-ish auto shape or shape named QR
            name_lower = shape.name.lower()
            if 'qr' in name_lower:
                qr_found = True
                break
            # Check if it's a roughly square auto shape
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                w = shape.width
                h = shape.height
                if w > 0 and h > 0:
                    ratio = max(w, h) / min(w, h)
                    if ratio <= 1.2 and w > 914400:  # roughly square and reasonably large (> 1 inch)
                        # Check if text mentions QR
                        if shape.has_text_frame:
                            txt = shape.text_frame.text.lower()
                            if 'qr' in txt:
                                qr_found = True
                                break

        if qr_found:
            print(f"PASS: Component 8 — Slide 10 has QR placeholder shape (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 8 — Slide 10 missing QR placeholder shape")
    except Exception as e:
        print(f"ERROR: Component 8 — {e}")

    # Component 9: Orange #E65100 accent color used across multiple slides (0.10 pts)
    try:
        slides_with_orange = 0
        for slide in prs.slides:
            colors = get_font_colors_from_slide(slide)
            if 'E65100' in colors:
                slides_with_orange += 1

        if slides_with_orange >= 5:
            print(f"PASS: Component 9 — Orange #E65100 found on {slides_with_orange}/10 slides (0.10 pts)")
            total_score += 0.10
        elif slides_with_orange >= 3:
            print(f"PARTIAL: Component 9 — Orange #E65100 found on {slides_with_orange} slides (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 9 — Orange #E65100 only on {slides_with_orange} slides, expected >=5")
    except Exception as e:
        print(f"ERROR: Component 9 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for unsaved GUI state
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Main entry point
if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(FILE_PATH)
