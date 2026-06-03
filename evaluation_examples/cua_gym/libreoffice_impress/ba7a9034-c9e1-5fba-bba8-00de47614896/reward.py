"""
Reward Script: Insert a 5-column, 6-row table on slide 4 with specific headers and first data row
Task ID: impress_tm_067
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.25): Table exists on slide 4 with 6 rows x 5 columns
  - Component 2 (0.25): Header row matches: Name, Role, Department, Start Date, Status
  - Component 3 (0.35): First data row matches: Alice Chen, Developer, Engineering, 2025-03-15, Active
  - Component 4 (0.15): Remaining rows (2-5) are empty
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_067'


def persist_app_state():
    """Save any unsaved LibreOffice Impress state."""
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        import time
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide4 = prs.slides[3]  # 0-indexed, so slide 4 is index 3

    # Find table on slide 4
    table_shape = None
    for shape in slide4.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_shape = shape
            break

    # Component 1: Table exists on slide 4 with correct dimensions (0.25 points)
    try:
        if table_shape is None:
            print("FAIL: Component 1 -- No table found on slide 4")
        else:
            table = table_shape.table
            num_rows = len(table.rows)
            num_cols = len(table.columns)
            if num_rows == 6 and num_cols == 5:
                print(f"PASS: Component 1 -- Table found on slide 4 with {num_rows} rows x {num_cols} cols (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 -- Table dimensions are {num_rows}x{num_cols}, expected 6x5")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # If no table found, remaining checks cannot proceed
    if table_shape is None:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    table = table_shape.table

    # Component 2: Header row (row 0) has correct values (0.25 points)
    try:
        expected_headers = ['Name', 'Role', 'Department', 'Start Date', 'Status']
        actual_headers = []
        for c in range(min(len(table.columns), 5)):
            actual_headers.append(table.cell(0, c).text.strip())

        if actual_headers == expected_headers:
            print(f"PASS: Component 2 -- Headers match: {actual_headers} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 -- Headers mismatch. Expected {expected_headers}, got {actual_headers}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: First data row (row 1) has correct values (0.35 points)
    try:
        expected_data = ['Alice Chen', 'Developer', 'Engineering', '2025-03-15', 'Active']
        actual_data = []
        for c in range(min(len(table.columns), 5)):
            actual_data.append(table.cell(1, c).text.strip())

        if actual_data == expected_data:
            print(f"PASS: Component 3 -- Data row 1 matches: {actual_data} (0.35 pts)")
            total_score += 0.35
        else:
            # Partial credit: count matching cells
            matching = sum(1 for a, e in zip(actual_data, expected_data) if a == e)
            if matching >= 3:
                partial = 0.35 * (matching / 5)
                print(f"PARTIAL: Component 3 -- {matching}/5 cells match. Expected {expected_data}, got {actual_data} ({partial:.2f} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 3 -- Data row 1 mismatch. Expected {expected_data}, got {actual_data}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Rows 2-5 are empty (0.15 points)
    try:
        all_empty = True
        non_empty_cells = []
        for r in range(2, min(len(table.rows), 6)):
            for c in range(min(len(table.columns), 5)):
                cell_text = table.cell(r, c).text.strip()
                if cell_text:
                    all_empty = False
                    non_empty_cells.append(f"({r},{c})={cell_text}")

        if all_empty:
            print(f"PASS: Component 4 -- Rows 2-5 are all empty (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 -- Non-empty cells in rows 2-5: {non_empty_cells}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
persist_app_state()

if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
