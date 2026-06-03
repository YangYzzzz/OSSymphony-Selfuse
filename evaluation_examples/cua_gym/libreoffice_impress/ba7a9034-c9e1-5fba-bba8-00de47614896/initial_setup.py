"""
Initial Setup: Insert a table on slide 4 of an HR overview presentation
Task ID: impress_tm_067
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_067'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "HR Overview"
    slide1.placeholders[1].text = "Annual Report 2025"

    # --- Slide 2: Employee Benefits ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Employee Benefits"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Comprehensive health insurance for all full-time employees"
    p = tf2.add_paragraph()
    p.text = "401(k) matching up to 6% of annual salary"
    p = tf2.add_paragraph()
    p.text = "Flexible remote work policy with 3 days per week option"
    p = tf2.add_paragraph()
    p.text = "Annual professional development budget of $2,500 per employee"
    p = tf2.add_paragraph()
    p.text = "Generous PTO: 20 days vacation + 10 sick days"

    # --- Slide 3: Performance Metrics ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Performance Metrics"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Q4 revenue growth: 18.3% year-over-year"
    p = tf3.add_paragraph()
    p.text = "Employee satisfaction score: 4.2 out of 5.0"
    p = tf3.add_paragraph()
    p.text = "Customer retention rate improved to 94.7%"
    p = tf3.add_paragraph()
    p.text = "Average project delivery time reduced by 12%"

    # --- Slide 4: Team Roster (TITLE ONLY - no table) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    tf4 = txBox.text_frame
    p = tf4.paragraphs[0]
    p.text = "Team Roster"
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 5: Training Schedule ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Training Schedule"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "March 20: Leadership Workshop - Building Effective Teams"
    p = tf5.add_paragraph()
    p.text = "April 5: Technical Deep Dive - Cloud Architecture Best Practices"
    p = tf5.add_paragraph()
    p.text = "April 18: Compliance Training - Data Privacy Regulations"
    p = tf5.add_paragraph()
    p.text = "May 10: Soft Skills - Presentation and Communication"

    # --- Slide 6: Budget Summary ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Budget Summary"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Total HR budget for FY2025: $3.2M"
    p = tf6.add_paragraph()
    p.text = "Recruitment costs: $480,000 (15% of total)"
    p = tf6.add_paragraph()
    p.text = "Training and development: $320,000 (10% of total)"
    p = tf6.add_paragraph()
    p.text = "Benefits administration: $1.6M (50% of total)"
    p = tf6.add_paragraph()
    p.text = "Remaining operational expenses: $800,000 (25% of total)"

    # --- Slide 7: Next Steps ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Next Steps"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "Finalize Q2 hiring plan by March 30"
    p = tf7.add_paragraph()
    p.text = "Roll out updated performance review framework in April"
    p = tf7.add_paragraph()
    p.text = "Launch employee engagement survey by mid-April"
    p = tf7.add_paragraph()
    p.text = "Present diversity metrics to leadership on May 1"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
