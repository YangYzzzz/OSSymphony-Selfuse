"""
Reward Script: Pharmaceutical Research Presentation (Drug_Pipeline.pptx)
Task ID: impress_wf_071
Domain: libreoffice_impress
Scoring:
  1. File on Desktop (0.10)
  2. 12 slides (0.15)
  3. Slide 1 title content (0.10)
  4. Slide 2 timeline with 5 phase markers (0.10)
  5. Slide 3 hexagonal shapes (0.10)
  6. Slide 5 preclinical table (0.05)
  7. Slide 7 bar chart (0.10)
  8. Slide 8 color-coded safety table (0.10)
  9. Slide 11 flowchart with decision shapes (0.10)
  10. Colors #0277BD and #4CAF50 used (0.10)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_071'
FILE_PATH = os.path.join(WORKDIR, 'Desktop', 'Drug_Pipeline.pptx')


def collect_all_colors(prs):
    """Collect all RGB color strings from shape fills and text runs across the presentation."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    colors = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            # Shape fill colors
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                try:
                    if shape.fill.type is not None:
                        c = str(shape.fill.fore_color.rgb)
                        colors.add(c.upper())
                except Exception:
                    pass
            # Text run colors
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color.type is not None:
                                c = str(run.font.color.rgb)
                                colors.add(c.upper())
                        except Exception:
                            pass
            # Table cell fill colors
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                for r in range(len(table.rows)):
                    for c in range(len(table.columns)):
                        try:
                            cell = table.cell(r, c)
                            if cell.fill.type is not None:
                                clr = str(cell.fill.fore_color.rgb)
                                colors.add(clr.upper())
                        except Exception:
                            pass
    # Also check via XML for thorough coverage
    try:
        with zipfile.ZipFile(FILE_PATH, 'r') as zf:
            for name in zf.namelist():
                if name.startswith('ppt/slides/slide') and name.endswith('.xml'):
                    root = ET.parse(zf.open(name)).getroot()
                    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                    for elem in root.findall('.//a:srgbClr', ns):
                        val = elem.get('val', '').upper()
                        if val:
                            colors.add(val)
    except Exception:
        pass
    return colors


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    total_score = 0.0

    # Precondition: load file
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: File exists on Desktop (0.10 points)
    # This is inherently checked since we loaded it, but we verify path specifically
    try:
        if os.path.exists(file_path) and 'Desktop' in file_path:
            print(f"PASS: Component 1 - Drug_Pipeline.pptx found on Desktop (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 1 - Drug_Pipeline.pptx not on Desktop")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Exactly 12 slides (0.15 points)
    try:
        if num_slides == 12:
            print(f"PASS: Component 2 - Exactly 12 slides found (0.15 pts)")
            total_score += 0.15
        elif num_slides >= 10:
            partial = 0.08
            print(f"PARTIAL: Component 2 - Found {num_slides} slides instead of 12 ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 - Found {num_slides} slides, expected 12")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Slide 1 has 'Drug Discovery Pipeline' and 'Compound XR-7' (0.10 points)
    try:
        if num_slides >= 1:
            slide1 = prs.slides[0]
            all_text = " ".join(
                shape.text_frame.text for shape in slide1.shapes if shape.has_text_frame
            ).lower()
            has_pipeline = "drug discovery pipeline" in all_text
            has_xr7 = "xr-7" in all_text or "xr7" in all_text or "compound xr" in all_text
            if has_pipeline and has_xr7:
                print(f"PASS: Component 3 - Slide 1 title has 'Drug Discovery Pipeline' and 'XR-7' (0.10 pts)")
                total_score += 0.10
            elif has_pipeline or has_xr7:
                print(f"PARTIAL: Component 3 - Slide 1 has only one of pipeline/XR-7 (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 3 - Slide 1 missing title content. Found: {all_text[:100]}")
        else:
            print(f"FAIL: Component 3 - No slides found")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Slide 2 has timeline with 5 phase markers (0.10 points)
    # Expected: Arrow shape + text containing Discovery, Preclinical, Phase I, Phase II, Phase III
    try:
        if num_slides >= 2:
            slide2 = prs.slides[1]
            all_text = " ".join(
                shape.text_frame.text for shape in slide2.shapes if shape.has_text_frame
            ).lower()
            phases = ["discovery", "preclinical", "phase i", "phase ii", "phase iii"]
            phases_found = sum(1 for p in phases if p in all_text)
            has_arrow = any(
                "arrow" in (shape.name or "").lower()
                for shape in slide2.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            )
            if phases_found >= 5 and has_arrow:
                print(f"PASS: Component 4 - Slide 2 has arrow + all 5 phases (0.10 pts)")
                total_score += 0.10
            elif phases_found >= 3:
                partial = 0.05
                print(f"PARTIAL: Component 4 - Found {phases_found}/5 phases, arrow={has_arrow} ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 4 - Slide 2 phases_found={phases_found}, arrow={has_arrow}")
        else:
            print(f"FAIL: Component 4 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Slide 3 has hexagonal shapes (0.10 points)
    try:
        if num_slides >= 3:
            slide3 = prs.slides[2]
            hexagon_count = sum(
                1 for shape in slide3.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                and "hexagon" in (shape.name or "").lower()
            )
            if hexagon_count >= 4:
                print(f"PASS: Component 5 - Slide 3 has {hexagon_count} hexagonal shapes (0.10 pts)")
                total_score += 0.10
            elif hexagon_count >= 2:
                print(f"PARTIAL: Component 5 - Slide 3 has {hexagon_count} hexagons (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 5 - Slide 3 has {hexagon_count} hexagons, expected >= 4")
        else:
            print(f"FAIL: Component 5 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    # Component 6: Slide 5 has a preclinical results table (0.05 points)
    try:
        if num_slides >= 5:
            slide5 = prs.slides[4]
            has_table = any(
                shape.shape_type == MSO_SHAPE_TYPE.TABLE for shape in slide5.shapes
            )
            if has_table:
                print(f"PASS: Component 6 - Slide 5 has a table (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 6 - Slide 5 has no table")
        else:
            print(f"FAIL: Component 6 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 6 - {e}")

    # Component 7: Slide 7 has a bar chart (0.10 points)
    try:
        if num_slides >= 7:
            slide7 = prs.slides[6]
            chart_shapes = [
                s for s in slide7.shapes if s.shape_type == MSO_SHAPE_TYPE.CHART
            ]
            if len(chart_shapes) > 0:
                ct = int(chart_shapes[0].chart.chart_type)
                # Accept column/bar chart types (51-54 are column variants, 57-60 are bar variants)
                if 51 <= ct <= 54 or 57 <= ct <= 60:
                    print(f"PASS: Component 7 - Slide 7 has bar/column chart type={ct} (0.10 pts)")
                    total_score += 0.10
                elif ct > 0:
                    print(f"PARTIAL: Component 7 - Slide 7 has chart but type={ct}, not bar/column (0.05 pts)")
                    total_score += 0.05
            else:
                print(f"FAIL: Component 7 - Slide 7 has no chart")
        else:
            print(f"FAIL: Component 7 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 7 - {e}")

    # Component 8: Slide 8 has color-coded safety table (0.10 points)
    # Table must exist AND have colored cells in the severity column
    try:
        if num_slides >= 8:
            slide8 = prs.slides[7]
            table_shapes = [s for s in slide8.shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE]
            colored_cells = 0
            if len(table_shapes) > 0:
                table = table_shapes[0].table
                for r in range(len(table.rows)):
                    for c in range(len(table.columns)):
                        try:
                            cell = table.cell(r, c)
                            if cell.fill.type is not None:
                                colored_cells += 1
                        except Exception:
                            pass
            if len(table_shapes) > 0 and colored_cells >= 3:
                print(f"PASS: Component 8 - Slide 8 table with {colored_cells} colored cells (0.10 pts)")
                total_score += 0.10
            elif len(table_shapes) > 0:
                print(f"PARTIAL: Component 8 - Slide 8 table found but only {colored_cells} colored cells (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 8 - Slide 8 has no table")
        else:
            print(f"FAIL: Component 8 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 8 - {e}")

    # Component 9: Slide 11 has flowchart with diamond (decision) shapes (0.10 points)
    try:
        if num_slides >= 11:
            slide11 = prs.slides[10]
            diamond_count = sum(
                1 for shape in slide11.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                and "diamond" in (shape.name or "").lower()
            )
            rounded_rect_count = sum(
                1 for shape in slide11.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                and "rounded" in (shape.name or "").lower()
            )
            if diamond_count >= 2 and rounded_rect_count >= 2:
                print(f"PASS: Component 9 - Slide 11 flowchart: {diamond_count} diamonds, {rounded_rect_count} rounded rects (0.10 pts)")
                total_score += 0.10
            elif diamond_count >= 1 or rounded_rect_count >= 2:
                print(f"PARTIAL: Component 9 - Slide 11: {diamond_count} diamonds, {rounded_rect_count} rounded rects (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 9 - Slide 11 missing decision/process shapes")
        else:
            print(f"FAIL: Component 9 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 9 - {e}")

    # Component 10: Colors #0277BD (medical blue) and #4CAF50 (green) used (0.10 points)
    try:
        all_colors = collect_all_colors(prs)
        has_blue = "0277BD" in all_colors
        has_green = "4CAF50" in all_colors
        if has_blue and has_green:
            print(f"PASS: Component 10 - Both #0277BD and #4CAF50 found (0.10 pts)")
            total_score += 0.10
        elif has_blue or has_green:
            found = "#0277BD" if has_blue else "#4CAF50"
            print(f"PARTIAL: Component 10 - Only {found} found (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 10 - Neither #0277BD nor #4CAF50 found. Colors: {all_colors}")
    except Exception as e:
        print(f"ERROR: Component 10 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(FILE_PATH)
