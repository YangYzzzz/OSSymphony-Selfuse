"""
reward.py for impress_ndo_080
Verifies that slides 2,3,4 from Source.pptx were copied into Target.pptx
after slide 5, preserving the blue source formatting.

Expected golden state:
- Target.pptx has 13 slides (was 10)
- Slides 6,7,8 are copies of Source slides 2,3,4 with blue theme
- Original Target slides 6-10 shifted to positions 9-13
"""

from pptx import Presentation

TARGET_PATH = "/home/user/Target.pptx"
SOURCE_PATH = "/home/user/Source.pptx"

# Blue theme background from Source
BLUE_BG = "0D2747"
# Green theme background from Target
GREEN_BG = "0B3D1E"

def get_slide_title(slide):
    """Extract the first non-empty text from a slide as its title."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    return t
    return ""

def get_slide_all_text(slide):
    """Get all text content from a slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return texts

def get_bg_color(slide):
    """Get background color as hex string or None."""
    try:
        fill = slide.background.fill
        if fill.type == 1:  # Solid fill
            return str(fill.fore_color.rgb)
        elif fill.type == 5:  # Inherited
            master_fill = slide.slide_layout.slide_master.background.fill
            if master_fill.type == 1:
                return str(master_fill.fore_color.rgb)
    except Exception:
        pass
    return None

def get_text_colors(slide):
    """Get all explicit text colors from a slide."""
    colors = set()
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.color.type is not None:
                            colors.add(str(run.font.color.rgb))
                    except Exception:
                        pass
    return colors


def main():
    score = 0.0

    try:
        tgt = Presentation(TARGET_PATH)
    except Exception as e:
        print(f"FAIL: Cannot open Target.pptx: {e}")
        print(f"REWARD: 0.0")
        return

    try:
        src = Presentation(SOURCE_PATH)
    except Exception as e:
        print(f"FAIL: Cannot open Source.pptx: {e}")
        print(f"REWARD: 0.0")
        return

    num_slides = len(tgt.slides)

    # ---- Component 1 (0.2): Target has exactly 13 slides ----
    if num_slides == 13:
        score += 0.2
        print("PASS Component 1: Target.pptx has 13 slides")
    else:
        print(f"FAIL Component 1: Target.pptx has {num_slides} slides (expected 13)")

    # ---- Component 2 (0.3): Slides 6,7,8 content matches Source slides 2,3,4 ----
    if num_slides >= 8:
        # Source slides 2,3,4 (0-indexed: 1,2,3)
        source_titles = [get_slide_title(src.slides[i]) for i in [1, 2, 3]]
        source_texts = [get_slide_all_text(src.slides[i]) for i in [1, 2, 3]]

        # Target slides 6,7,8 (0-indexed: 5,6,7)
        target_titles = [get_slide_title(tgt.slides[i]) for i in [5, 6, 7]]
        target_texts = [get_slide_all_text(tgt.slides[i]) for i in [5, 6, 7]]

        matches = 0
        for idx in range(3):
            if source_titles[idx] and source_titles[idx] == target_titles[idx]:
                # Also check that body text overlaps
                src_set = set(source_texts[idx])
                tgt_set = set(target_texts[idx])
                if src_set and src_set.issubset(tgt_set):
                    matches += 1
                    print(f"  PASS: Slide {idx+6} title matches Source slide {idx+2}: '{source_titles[idx]}'")
                else:
                    print(f"  PARTIAL: Slide {idx+6} title matches but body text differs")
                    matches += 0.5
            else:
                print(f"  FAIL: Slide {idx+6} title='{target_titles[idx]}' != Source slide {idx+2} title='{source_titles[idx]}'")

        comp2_score = (matches / 3.0) * 0.3
        score += comp2_score
        print(f"PASS Component 2: Content match score {matches}/3 ({comp2_score:.2f})")
    else:
        print("FAIL Component 2: Not enough slides to check content")

    # ---- Component 3 (0.3): Slides 6,7,8 have blue theme (not green) ----
    if num_slides >= 8:
        blue_count = 0
        for idx in [5, 6, 7]:
            bg = get_bg_color(tgt.slides[idx])
            text_colors = get_text_colors(tgt.slides[idx])
            # Check background is blue
            if bg == BLUE_BG:
                blue_count += 1
                print(f"  PASS: Slide {idx+1} has blue background ({bg})")
            elif bg == GREEN_BG:
                print(f"  FAIL: Slide {idx+1} has green background ({bg}), expected blue ({BLUE_BG})")
            else:
                # Check if text colors indicate blue theme (CBDBEF is the blue accent)
                if "CBDBEF" in text_colors:
                    blue_count += 0.5
                    print(f"  PARTIAL: Slide {idx+1} bg={bg} but has blue theme text colors")
                else:
                    print(f"  FAIL: Slide {idx+1} bg={bg}, no blue theme indicators")

        comp3_score = (blue_count / 3.0) * 0.3
        score += comp3_score
        print(f"PASS Component 3: Blue theme score {blue_count}/3 ({comp3_score:.2f})")
    else:
        print("FAIL Component 3: Not enough slides to check theme")

    # ---- Component 4 (0.2): Original slides 6-10 preserved at positions 9-13 ----
    if num_slides >= 13:
        # The original Target slide titles at positions 6-10 (0-indexed 5-9)
        # These should now be at positions 9-13 (0-indexed 8-12)
        expected_titles = [
            "Community Engagement Programs",
            "Green Building Certifications",
            "Biodiversity Conservation Efforts",
            "Stakeholder Transparency Report",
            "2026 Sustainability Roadmap",
        ]
        preserved = 0
        for i, expected in enumerate(expected_titles):
            actual = get_slide_title(tgt.slides[8 + i])
            if actual == expected:
                preserved += 1
                print(f"  PASS: Slide {9+i} title='{actual}'")
            else:
                print(f"  FAIL: Slide {9+i} title='{actual}' != expected '{expected}'")

        comp4_score = (preserved / 5.0) * 0.2
        score += comp4_score
        print(f"PASS Component 4: Preserved slides {preserved}/5 ({comp4_score:.2f})")
    else:
        print("FAIL Component 4: Not enough slides to check preservation")

    # Round to 1 decimal
    score = round(score, 1)
    # Clamp
    score = max(0.0, min(1.0, score))
    print(f"REWARD: {score}")

if __name__ == "__main__":
    main()
