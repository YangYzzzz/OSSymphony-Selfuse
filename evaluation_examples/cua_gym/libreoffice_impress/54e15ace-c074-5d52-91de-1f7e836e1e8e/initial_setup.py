"""
initial_setup.py — Create Source.pptx (8 slides, blue theme) and Target.pptx (10 slides, green theme).
Opens both files in LibreOffice Impress.
"""

import os
import subprocess
import shlex
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── Color palettes ──────────────────────────────────────────────────────────
BLUE_BG = RGBColor(0x0D, 0x27, 0x47)        # dark navy background
BLUE_ACCENT = RGBColor(0x1E, 0x90, 0xFF)     # dodger blue accent
BLUE_TITLE = RGBColor(0xFF, 0xFF, 0xFF)       # white title text
BLUE_BODY = RGBColor(0xCB, 0xDB, 0xEF)       # light blue body text

GREEN_BG = RGBColor(0x0B, 0x3D, 0x1E)        # dark forest background
GREEN_ACCENT = RGBColor(0x2E, 0xCC, 0x71)    # emerald accent
GREEN_TITLE = RGBColor(0xFF, 0xFF, 0xFF)      # white title text
GREEN_BODY = RGBColor(0xC8, 0xF7, 0xD5)      # light green body text


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_textbox(slide, text, left, top, width, height, font_size, color, bold=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color


def add_body_textbox(slide, text, left, top, width, height, font_size, color):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = font_size
    run.font.bold = False
    run.font.color.rgb = color


def add_accent_line(slide, color):
    """Add a colored accent line near the top of the slide."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0.5), Inches(1.6), Inches(1.5), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def create_blue_slide(prs, title, body):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BLUE_BG)
    add_accent_line(slide, BLUE_ACCENT)
    add_title_textbox(slide, title, Inches(0.5), Inches(0.4), Inches(9), Inches(1.2), Pt(32), BLUE_TITLE)
    add_body_textbox(slide, body, Inches(0.5), Inches(1.9), Inches(9), Inches(4.5), Pt(18), BLUE_BODY)
    return slide


def create_green_slide(prs, title, body):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, GREEN_BG)
    add_accent_line(slide, GREEN_ACCENT)
    add_title_textbox(slide, title, Inches(0.5), Inches(0.4), Inches(9), Inches(1.2), Pt(32), GREEN_TITLE)
    add_body_textbox(slide, body, Inches(0.5), Inches(1.9), Inches(9), Inches(4.5), Pt(18), GREEN_BODY)
    return slide


# ── Source.pptx — 8 slides, blue theme ─────────────────────────────────────
source_slides = [
    ("Q3 2025 Strategic Review", "This presentation covers the key strategic initiatives, financial performance, and forward-looking plans for the third quarter of fiscal year 2025."),
    ("Revenue Performance Overview", "Total revenue reached $48.2M in Q3, representing a 17% year-over-year increase. Subscription revenue grew 23% while professional services revenue remained stable at $8.1M."),
    ("Customer Acquisition Metrics", "We acquired 1,240 new enterprise customers this quarter, bringing total active accounts to 8,750. Customer acquisition cost decreased by 12% due to improved marketing efficiency."),
    ("Product Development Roadmap", "Three major product releases shipped on schedule. The AI-powered analytics module achieved 94% accuracy in beta testing. Mobile app redesign reduced load times by 40%."),
    ("Market Expansion Strategy", "Entered two new geographic markets: Southeast Asia and Northern Europe. Partnership agreements signed with 15 regional distributors. Localization completed for 8 additional languages."),
    ("Operational Efficiency Gains", "Cloud infrastructure costs reduced by 22% through optimization. Automated deployment pipeline cut release cycles from 2 weeks to 3 days. Support ticket resolution time improved by 35%."),
    ("Talent & Workforce Update", "Headcount grew to 620 employees across 12 offices. Engineering team expanded by 45 new hires. Employee satisfaction score reached 4.3 out of 5.0 in quarterly survey."),
    ("Q4 Outlook and Key Priorities", "Projected Q4 revenue of $52M-$55M. Priority initiatives include enterprise platform launch, Series C fundraising preparation, and expansion of the partner ecosystem."),
]

src_prs = Presentation()
for title, body in source_slides:
    create_blue_slide(src_prs, title, body)

src_prs.save("/home/user/Source.pptx")
print("Created Source.pptx with 8 slides (blue theme)")


# ── Target.pptx — 10 slides, green theme ───────────────────────────────────
target_slides = [
    ("Sustainability Annual Report 2025", "This report outlines our environmental commitments, progress on sustainability goals, and the roadmap for achieving carbon neutrality by 2030."),
    ("Environmental Impact Summary", "Total carbon emissions reduced by 28% compared to 2023 baseline. Renewable energy now powers 72% of our global operations. Water usage decreased by 15% across all facilities."),
    ("Renewable Energy Transition", "Solar installations completed at 8 additional facilities this year. Wind power purchase agreements signed for 120 MW capacity. Battery storage systems deployed at 3 data centers."),
    ("Waste Reduction Initiatives", "Achieved 89% waste diversion rate, up from 76% last year. Single-use plastics eliminated from all office locations. Packaging redesign reduced material usage by 34%."),
    ("Supply Chain Sustainability", "85% of tier-one suppliers now meet our sustainability standards. Blockchain-based tracking implemented for raw material sourcing. Transportation emissions cut by 19% through route optimization."),
    ("Community Engagement Programs", "Invested $3.2M in local environmental restoration projects. Employee volunteer hours reached 42,000 across 150 community events. Launched STEM education partnership with 25 schools."),
    ("Green Building Certifications", "Four additional offices achieved LEED Platinum certification. Average energy consumption per square foot decreased by 21%. Smart building systems installed in 90% of owned properties."),
    ("Biodiversity Conservation Efforts", "Protected 5,000 acres of natural habitat through land trust partnerships. Corporate campus wildlife corridors expanded. Native plant restoration completed at 12 facility locations."),
    ("Stakeholder Transparency Report", "Published comprehensive ESG data for the third consecutive year. Third-party audit confirmed 98% data accuracy. Investor sustainability rating improved from B+ to A-."),
    ("2026 Sustainability Roadmap", "Target 50% absolute emission reduction by end of 2026. Plan to achieve 100% renewable energy for operations. Launch circular economy pilot program across product lines."),
]

tgt_prs = Presentation()
for title, body in target_slides:
    create_green_slide(tgt_prs, title, body)

tgt_prs.save("/home/user/Target.pptx")
print("Created Target.pptx with 10 slides (green theme)")


# ── Open both in LibreOffice Impress ────────────────────────────────────────
def launch_gui(command, delay_sec=2.0):
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)

launch_gui('libreoffice --impress "/home/user/Source.pptx"', delay_sec=3.0)
launch_gui('libreoffice --impress "/home/user/Target.pptx"', delay_sec=2.0)

print("Both files opened in LibreOffice Impress.")
