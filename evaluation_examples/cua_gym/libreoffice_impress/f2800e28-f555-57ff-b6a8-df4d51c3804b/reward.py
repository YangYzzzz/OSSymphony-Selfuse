"""
Reward Script: Insert 4 images in a 2x2 grid layout on slide 4
Task ID: impress_stu_047
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): Slide 4 has exactly 4 pictures
  Component 2 (0.3): All pictures are approximately 3x3 inches
  Component 3 (0.4): Pictures are arranged in a 2x2 grid (2 columns, 2 rows, even spacing)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_047'
EMU_PER_INCH = 914400


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide4 = prs.slides[3]  # 0-indexed, slide 4

    # Collect all picture shapes on slide 4
    pictures = []
    for shape in slide4.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pictures.append(shape)

    # Component 1: Slide 4 has exactly 4 pictures (0.3 points)
    try:
        num_pics = len(pictures)
        if num_pics == 4:
            print(f"PASS: Component 1 -- Slide 4 has exactly 4 pictures (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 -- Expected 4 pictures on slide 4, found {num_pics}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # If not exactly 4 pictures, remaining checks cannot meaningfully pass
    if len(pictures) != 4:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: All pictures are approximately 3x3 inches (0.3 points)
    # Tolerance: 15% of 3 inches = 0.45 inches
    try:
        target_size_emu = 3 * EMU_PER_INCH  # 2743200 EMU
        tolerance_emu = 0.45 * EMU_PER_INCH  # 411480 EMU
        correct_size_count = 0
        for i, pic in enumerate(pictures):
            w_ok = abs(pic.width - target_size_emu) <= tolerance_emu
            h_ok = abs(pic.height - target_size_emu) <= tolerance_emu
            w_in = pic.width / EMU_PER_INCH
            h_in = pic.height / EMU_PER_INCH
            if w_ok and h_ok:
                correct_size_count += 1
                print(f"  Picture {i+1} size OK: {w_in:.2f}x{h_in:.2f}in")
            else:
                print(f"FAIL: Component 2 -- Picture {i+1} size {w_in:.2f}x{h_in:.2f}in, expected ~3x3in")
        if correct_size_count == 4:
            print(f"PASS: Component 2 -- All 4 pictures are approximately 3x3 inches (0.3 pts)")
            total_score += 0.3
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Pictures are in a 2x2 grid layout (0.4 points)
    # Verify: 2 distinct left positions (columns) and 2 distinct top positions (rows),
    # with 2 pictures in each row and 2 in each column, and even spacing.
    try:
        lefts = sorted(set(pic.left for pic in pictures))
        tops = sorted(set(pic.top for pic in pictures))

        grid_issues = 0

        # Must have exactly 2 distinct columns and 2 distinct rows
        if len(lefts) != 2:
            print(f"FAIL: Component 3 -- Expected 2 distinct column positions, found {len(lefts)}: {[l/EMU_PER_INCH for l in lefts]}")
            grid_issues += 1
        if len(tops) != 2:
            print(f"FAIL: Component 3 -- Expected 2 distinct row positions, found {len(tops)}: {[t/EMU_PER_INCH for t in tops]}")
            grid_issues += 1

        if grid_issues == 0:
            # Each (left, top) combination should have exactly 1 picture
            grid_cells = {}
            for pic in pictures:
                col = 0 if pic.left == lefts[0] else 1
                row = 0 if pic.top == tops[0] else 1
                cell = (row, col)
                if cell in grid_cells:
                    print(f"FAIL: Component 3 -- Multiple pictures in grid cell ({row},{col})")
                    grid_issues += 1
                else:
                    grid_cells[cell] = pic

            if grid_issues == 0 and len(grid_cells) == 4:
                # Verify columns are spaced apart (gap between them should be > 0)
                col_gap = lefts[1] - lefts[0]
                row_gap = tops[1] - tops[0]
                col_gap_in = col_gap / EMU_PER_INCH
                row_gap_in = row_gap / EMU_PER_INCH
                print(f"  Column spacing: {col_gap_in:.2f}in, Row spacing: {row_gap_in:.2f}in")

                # The columns should be spaced significantly (at least 2 inches apart center-to-center)
                if col_gap < 2 * EMU_PER_INCH:
                    print(f"FAIL: Component 3 -- Columns too close: {col_gap_in:.2f}in apart")
                    grid_issues += 1
                if row_gap < 2 * EMU_PER_INCH:
                    print(f"FAIL: Component 3 -- Rows too close: {row_gap_in:.2f}in apart")
                    grid_issues += 1

        if grid_issues == 0:
            print(f"PASS: Component 3 -- 2x2 grid layout verified with even spacing (0.4 pts)")
            total_score += 0.4
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
