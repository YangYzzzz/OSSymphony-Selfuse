"""
Initial Setup: Create a Microscopy Lab presentation with 6 slides and sample images.
Task ID: impress_stu_047
Domain: libreoffice_impress
Slide 4 is titled 'Microscope Observations' and is intentionally left empty (no images).
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_047'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
DOWNLOADS = f'{WORKDIR}/Downloads'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_sample_images():
    """Create 4 distinct sample microscopy-style images."""
    os.makedirs(DOWNLOADS, exist_ok=True)

    # Each image is a unique solid color with a simple pattern to simulate microscopy
    colors_and_labels = [
        ((180, 200, 220), (50, 80, 120), "Sample A - Epithelial Cells"),
        ((200, 220, 180), (80, 120, 50), "Sample B - Blood Smear"),
        ((220, 190, 200), (120, 50, 80), "Sample C - Muscle Tissue"),
        ((190, 210, 220), (50, 100, 110), "Sample D - Nerve Fibers"),
    ]

    for i, (bg_color, fg_color, label) in enumerate(colors_and_labels, 1):
        img = Image.new('RGB', (600, 600), bg_color)
        pixels = img.load()
        # Draw a simple cross-hatch pattern to mimic microscopy grid
        for x in range(600):
            for y in range(600):
                if x % 60 == 0 or y % 60 == 0:
                    pixels[x, y] = fg_color
                # Add some circular features
                cx, cy = 300, 300
                dist = ((x - cx) ** 2 + (y - cy) ** 2) ** 0.5
                if 140 < dist < 145 or 200 < dist < 205:
                    pixels[x, y] = fg_color

        path = f'{DOWNLOADS}/sample{i}.png'
        img.save(path)
        print(f'Created sample image: {path}')


def create_presentation():
    prs = Presentation()
    # Standard 10x7.5 widescreen
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Microscopy Lab Report"
    slide1.placeholders[1].text = "Biology 301 — Spring 2025\nDr. Elena Martinez"

    # --- Slide 2: Introduction to Microscopy ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Introduction to Microscopy"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Optical microscopy is one of the oldest scientific instruments"
    p2 = body2.add_paragraph()
    p2.text = "Modern compound microscopes achieve magnifications of 1000x or more"
    p2.level = 0
    p3 = body2.add_paragraph()
    p3.text = "Resolution is limited by the wavelength of visible light (~200 nm)"
    p3.level = 0
    p4 = body2.add_paragraph()
    p4.text = "Phase contrast and fluorescence techniques extend capabilities"
    p4.level = 0

    # --- Slide 3: Equipment & Materials ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Equipment & Materials"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Olympus BX53 compound microscope (4x, 10x, 40x, 100x objectives)"
    items3 = [
        "Prepared glass slides and coverslips",
        "Immersion oil (Type B, n=1.515)",
        "Methylene blue and H&E staining kits",
        "Digital camera attachment (Olympus DP74)",
        "Stage micrometer for calibration",
    ]
    for item in items3:
        p = body3.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 4: Microscope Observations (EMPTY - no images) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Microscope Observations"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # --- Slide 5: Data Analysis ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Data Analysis"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Cell counts were performed across 5 fields of view per sample"
    items5 = [
        "Average epithelial cell diameter: 12.4 ± 1.8 μm",
        "Red blood cell count: 4.7 × 10⁶ cells/μL (within normal range)",
        "Muscle fiber striations visible at 40x magnification",
        "Nerve fiber myelination confirmed with osmium tetroxide staining",
        "All measurements calibrated against 10 μm stage micrometer",
    ]
    for item in items5:
        p = body5.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 6: Conclusions ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Conclusions"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "All four tissue samples were successfully imaged and characterized"
    items6 = [
        "Epithelial cells showed expected morphology and size distribution",
        "Blood smear analysis confirmed normal erythrocyte parameters",
        "Skeletal muscle samples displayed clear banding patterns",
        "Peripheral nerve cross-sections revealed intact myelin sheaths",
        "Future work: electron microscopy for ultrastructural analysis",
    ]
    for item in items6:
        p = body6.add_paragraph()
        p.text = item
        p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')


def main():
    create_sample_images()
    create_presentation()

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


main()
