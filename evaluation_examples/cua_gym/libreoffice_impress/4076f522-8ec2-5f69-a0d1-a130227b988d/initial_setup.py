"""
Initial Setup: Create a 12-slide Year in Review presentation with empty slide 8
Task ID: impress_exec_019
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_019'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
    return slide


def add_title_only_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(prs, "FY2025 Year in Review", "Prepared by the Strategy & Operations Team\nQ4 Board Presentation")

    # Slide 2: Agenda
    add_content_slide(prs, "Agenda", [
        "Executive Summary",
        "Revenue & Financial Highlights",
        "Product Development Milestones",
        "Customer Acquisition & Retention",
        "Team Growth & Culture",
        "Year-over-Year Comparison",
        "Strategic Outlook FY2026",
        "Q&A"
    ])

    # Slide 3: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "FY2025 marked a transformational year for the organization",
        "Revenue grew 30.2% year-over-year to $62.5M",
        "Successfully launched 3 new product lines in enterprise segment",
        "Expanded into 4 new international markets (UK, DE, JP, AU)",
        "Net Promoter Score improved from 72 to 81"
    ])

    # Slide 4: Revenue Highlights
    add_content_slide(prs, "Revenue & Financial Highlights", [
        "Total Revenue: $62.5M (up from $48M in FY2024)",
        "Recurring Revenue: $41.3M (66% of total, up from 58%)",
        "Gross Margin: 61% (up from 58%)",
        "Net Income: $8.1M (up from $5.2M)",
        "Cash Position: $24.7M with zero long-term debt"
    ])

    # Slide 5: Product Milestones
    add_content_slide(prs, "Product Development Milestones", [
        "Enterprise Analytics Suite launched in Q1 - 47 customers onboarded",
        "Mobile app redesign completed with 4.7 App Store rating",
        "API v3 released with 99.97% uptime SLA",
        "AI-powered recommendations engine deployed in Q3",
        "Security certification: SOC 2 Type II achieved"
    ])

    # Slide 6: Customer Metrics
    add_content_slide(prs, "Customer Acquisition & Retention", [
        "Total Active Customers: 1,247 (up 34% YoY)",
        "Enterprise Accounts: 89 (up from 52)",
        "Customer Retention Rate: 94.3%",
        "Average Contract Value: $50.1K (up 18%)",
        "Support Ticket Resolution: < 4 hours average"
    ])

    # Slide 7: Team Growth
    add_content_slide(prs, "Team Growth & Culture", [
        "Total Headcount: 312 (up from 245 in FY2024)",
        "Engineering: 142 (+38 new hires)",
        "Sales & Marketing: 67 (+15 new hires)",
        "Employee Satisfaction Score: 4.6/5.0",
        "Voluntary Turnover: 8.2% (industry avg: 13.5%)"
    ])

    # Slide 8: Year-over-Year Comparison (TITLE ONLY - NO TABLE)
    add_title_only_slide(prs, "Year-over-Year Comparison")

    # Slide 9: Market Expansion
    add_content_slide(prs, "International Market Expansion", [
        "UK Office: London - 18 employees, $3.2M ARR",
        "Germany Office: Berlin - 12 employees, $2.1M ARR",
        "Japan Office: Tokyo - 8 employees, $1.4M ARR",
        "Australia Office: Sydney - 6 employees, $0.9M ARR",
        "International revenue now 12% of total (target: 20% by FY2027)"
    ])

    # Slide 10: Strategic Initiatives
    add_content_slide(prs, "FY2026 Strategic Initiatives", [
        "Launch AI Copilot for enterprise tier customers",
        "Expand APAC presence with Singapore hub",
        "Achieve $100M ARR milestone by Q4 FY2026",
        "Build partner ecosystem with 50+ integrations",
        "Invest $8M in R&D for next-gen platform"
    ])

    # Slide 11: Financial Outlook
    add_content_slide(prs, "FY2026 Financial Outlook", [
        "Revenue Target: $85M (36% growth)",
        "Gross Margin Target: 63%",
        "Planned Headcount: 420 employees",
        "Capital Expenditure: $12M",
        "Expected Net Income: $11.5M"
    ])

    # Slide 12: Q&A
    add_title_slide(prs, "Questions & Discussion", "Thank you for your continued support\nContact: strategy@company.com")

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
