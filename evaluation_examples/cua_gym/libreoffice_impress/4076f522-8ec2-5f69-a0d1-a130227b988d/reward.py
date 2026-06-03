"""
Reward Script: Create comparison table on slide 8
Task ID: impress_exec_019
Domain: libreoffice_impress
Scoring:
  - Component 1: Table exists with correct dimensions (0.2)
  - Component 2: Header row content (0.15)
  - Component 3: Data cell content (0.35)
  - Component 4: Alternating row shading (0.3)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_019'

# Expected table data: row 0 = header, rows 1-5 = data
EXPECTED_DATA = [
    ['Metric', 'FY2024', 'FY2025'],
    ['Revenue', '$48M', '$62.5M'],
    ['Gross Margin', '58%', '61%'],
    ['Net Income', '$5.2M', '$8.1M'],
    ['Headcount', '245', '312'],
    ['NPS Score', '72', '81'],
]

# Expected alternating fill: even rows white (FFFFFF), odd data rows light blue (E8F0FE)
# Row 0 (header): white, Row 1: white, Row 2: E8F0FE, Row 3: white, Row 4: E8F0FE, Row 5: white
EXPECTED_FILLS = {
    0: 'FFFFFF',
    1: 'FFFFFF',
    2: 'E8F0FE',
    3: 'FFFFFF',
    4: 'E8F0FE',
    5: 'FFFFFF',
}


def find_table_on_slide(slide):
    """Find a TABLE shape on the given slide. Returns table or None."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return shape.table
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 8 slides
    if len(prs.slides) < 8:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 8")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[7]  # 0-indexed, slide 8

    # Component 1: Table exists on slide 8 with 6 rows x 3 columns (0.2 points)
    try:
        table = find_table_on_slide(slide)
        if table is None:
            print("FAIL: Component 1 -- No table found on slide 8")
            print("REWARD: 0.0")
            return 0.0  # No table means nothing else can be checked

        num_rows = len(table.rows)
        num_cols = len(table.columns)
        if num_rows == 6 and num_cols == 3:
            print(f"PASS: Component 1 -- Table found with {num_rows}x{num_cols} (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 1 -- Table is {num_rows}x{num_cols}, expected 6x3")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Header row content matches (0.15 points)
    try:
        if table is not None and len(table.rows) >= 1 and len(table.columns) >= 3:
            header_match = True
            for c in range(3):
                actual = table.cell(0, c).text.strip()
                expected = EXPECTED_DATA[0][c]
                if actual != expected:
                    print(f"FAIL: Component 2 -- Header cell(0,{c}): expected {repr(expected)}, got {repr(actual)}")
                    header_match = False
            if header_match:
                print(f"PASS: Component 2 -- Header row matches: Metric, FY2024, FY2025 (0.15 pts)")
                total_score += 0.15
        else:
            print("FAIL: Component 2 -- Table dimensions insufficient for header check")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Data cell content matches (0.35 points, 0.07 per row)
    try:
        if table is not None and len(table.rows) >= 6 and len(table.columns) >= 3:
            data_score = 0.0
            points_per_row = 0.07
            for r in range(1, 6):
                row_match = True
                for c in range(3):
                    actual = table.cell(r, c).text.strip()
                    expected = EXPECTED_DATA[r][c]
                    if actual != expected:
                        print(f"FAIL: Component 3 -- cell({r},{c}): expected {repr(expected)}, got {repr(actual)}")
                        row_match = False
                if row_match:
                    data_score += points_per_row
                    print(f"PASS: Component 3 -- Row {r} data matches ({points_per_row} pts)")
            total_score += data_score
            print(f"Component 3 subtotal: {data_score:.2f}/0.35")
        else:
            print("FAIL: Component 3 -- Table dimensions insufficient for data check")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Alternating row shading (0.3 points, 0.05 per row)
    try:
        if table is not None and len(table.rows) >= 6 and len(table.columns) >= 3:
            shade_score = 0.0
            points_per_row = 0.05
            for r in range(6):
                expected_color = EXPECTED_FILLS[r]
                # Check fill of first cell in each row as representative
                cell = table.cell(r, 0)
                try:
                    fill = cell.fill
                    if fill.type is not None and fill.type == 1:  # SOLID
                        actual_color = str(fill.fore_color.rgb).upper()
                        if actual_color == expected_color:
                            shade_score += points_per_row
                            print(f"PASS: Component 4 -- Row {r} fill={actual_color} matches expected {expected_color} ({points_per_row} pts)")
                        else:
                            print(f"FAIL: Component 4 -- Row {r} fill={actual_color}, expected {expected_color}")
                    else:
                        print(f"FAIL: Component 4 -- Row {r} fill type is {fill.type}, expected SOLID (1)")
                except Exception as e:
                    print(f"FAIL: Component 4 -- Row {r} fill error: {e}")
            total_score += shade_score
            print(f"Component 4 subtotal: {shade_score:.2f}/0.30")
        else:
            print("FAIL: Component 4 -- Table dimensions insufficient for shading check")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
