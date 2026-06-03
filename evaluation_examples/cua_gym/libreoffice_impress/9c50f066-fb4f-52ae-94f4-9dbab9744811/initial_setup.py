"""
Initial Setup: Create a 20-slide presentation with varying quality issues for audit
Task ID: impress_gf5_043
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_043'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


# ---- Slide Content Design ----
# We'll create 20 slides with intentional quality issues:
#
# NO speaker notes: slides 2, 5, 8, 11, 14, 17  (6 slides)
# Overloaded (>100 words body): slides 3, 7, 12, 16  (4 slides)
# No images: slides 1, 4, 6, 9, 10, 13, 15, 18  (8 slides)
# Title > 60 chars: slides 4, 10, 19  (3 slides)
#
# Slides with images: 2, 3, 5, 7, 8, 11, 12, 14, 16, 17, 19, 20

SLIDE_DATA = [
    {  # Slide 1
        "title": "Q3 2025 Strategic Review",
        "body": "Welcome to the quarterly strategic review. This presentation covers our performance across all departments, key initiatives, and forward-looking plans for the remainder of the fiscal year.",
        "notes": "Opening slide - greet stakeholders, mention agenda timing is 45 minutes.",
        "has_image": False,
    },
    {  # Slide 2
        "title": "Revenue Performance Overview",
        "body": "Total revenue reached $12.4M in Q3, representing a 15% year-over-year increase. North American markets contributed 62% while EMEA grew by 23%.",
        "notes": "",  # NO notes
        "has_image": True,
    },
    {  # Slide 3 - overloaded
        "title": "Market Analysis Deep Dive",
        "body": (
            "Our comprehensive market analysis reveals several critical trends that demand immediate attention from the leadership team. "
            "First, the competitive landscape has shifted dramatically with three new entrants in the enterprise segment capturing approximately "
            "eight percent of market share within their first two quarters of operation. Second, customer acquisition costs have risen by "
            "seventeen percent across all channels, with digital advertising showing the steepest increase at twenty-two percent. Third, "
            "our retention rates remain strong at ninety-one percent, which is four percentage points above the industry average. However, "
            "the net promoter score has declined from seventy-eight to seventy-two, suggesting potential churn risks in upcoming quarters. "
            "Additionally, our enterprise pipeline has grown by thirty-five percent, driven largely by expansion in the healthcare and "
            "financial services verticals. The mid-market segment shows promise with several pilot programs converting to full deployments. "
            "We recommend increasing investment in product differentiation and customer success programs to maintain our competitive position."
        ),
        "notes": "Spend extra time on competitive analysis. Board will ask about pricing strategy.",
        "has_image": True,
    },
    {  # Slide 4 - long title, no image
        "title": "Comprehensive Breakdown of Regional Sales Performance Across All Territories and Product Lines",
        "body": "Northeast: $3.2M (+18%), Southeast: $2.1M (+12%), Midwest: $1.8M (+8%), West Coast: $3.5M (+22%), International: $1.8M (+31%).",
        "notes": "Highlight West Coast growth driven by new partnership with TechVantage Solutions.",
        "has_image": False,
    },
    {  # Slide 5
        "title": "Product Development Milestones",
        "body": "Platform v4.2 launched on schedule. New AI-powered analytics module achieved 98.5% accuracy in beta testing with 200 enterprise users.",
        "notes": "",  # NO notes
        "has_image": True,
    },
    {  # Slide 6
        "title": "Engineering Team Updates",
        "body": "Hired 12 senior engineers in Q3. Reduced deployment cycle time from 4 weeks to 10 days. Achieved 99.97% uptime SLA compliance.",
        "notes": "Mention the DevOps transformation initiative led by Priya Sharma's team.",
        "has_image": False,
    },
    {  # Slide 7 - overloaded
        "title": "Customer Success Stories",
        "body": (
            "Our customer success team has documented remarkable outcomes across multiple verticals this quarter. Meridian Healthcare "
            "reduced patient intake processing time by forty-five percent after implementing our workflow automation suite, resulting in "
            "annual savings of approximately one point two million dollars. Summit Financial Group deployed our risk analytics platform "
            "across their investment division, achieving a thirty percent improvement in portfolio risk assessment accuracy. NovaTech "
            "Manufacturing integrated our supply chain optimization module, which identified over two million dollars in cost reduction "
            "opportunities within the first sixty days. Furthermore, Pacific Retail Group leveraged our customer analytics engine to "
            "increase online conversion rates by twenty-eight percent during their peak holiday season. These success stories demonstrate "
            "the tangible business value our platform delivers across diverse industry verticals and use cases. The customer success "
            "team recommends expanding case study development to support sales enablement and marketing content creation efforts."
        ),
        "notes": "Use Meridian Healthcare as the lead case study. CEO Sarah Kim available for testimonial.",
        "has_image": True,
    },
    {  # Slide 8
        "title": "Marketing Campaign Results",
        "body": "Digital campaigns generated 4,200 qualified leads. Content marketing drove 340K unique visitors. Brand awareness increased 18% in target demographics.",
        "notes": "",  # NO notes
        "has_image": True,
    },
    {  # Slide 9
        "title": "Financial Summary",
        "body": "Gross margin: 72.3%, Operating expenses: $8.9M, EBITDA: $3.5M, Cash position: $28.4M, Runway: 36+ months.",
        "notes": "CFO will present detailed P&L in the appendix section if requested.",
        "has_image": False,
    },
    {  # Slide 10 - long title, no image
        "title": "Year-over-Year Comparison of Key Performance Indicators Across All Business Units and Departments",
        "body": "Revenue: +15%, Customers: +22%, NPS: -6pts, Employee satisfaction: +8pts, Time-to-market: -35%, Support tickets: -12%.",
        "notes": "NPS decline needs explanation. Primarily driven by onboarding complexity in v4.0 release.",
        "has_image": False,
    },
    {  # Slide 11
        "title": "Partnership Ecosystem",
        "body": "Signed 8 new technology partners. Strategic alliance with CloudBridge expanded to cover 14 countries. Channel revenue up 28%.",
        "notes": "",  # NO notes
        "has_image": True,
    },
    {  # Slide 12 - overloaded
        "title": "Talent Acquisition and Retention",
        "body": (
            "The human resources division achieved significant milestones in talent acquisition and retention during the third quarter. "
            "Total headcount grew from three hundred twelve to three hundred forty-eight employees, representing a net increase of "
            "thirty-six positions across engineering, sales, and customer success departments. The voluntary attrition rate decreased "
            "to seven point two percent, down from nine point eight percent in the previous quarter, attributed to the implementation "
            "of our enhanced benefits package and flexible remote work policy. The learning and development team launched fifteen new "
            "training programs, with an average completion rate of eighty-seven percent. Employee engagement survey scores improved "
            "across all categories, with notable gains in career development satisfaction rising from sixty-two to seventy-four percent. "
            "Diversity hiring initiatives resulted in forty-two percent of new hires identifying as underrepresented minorities. The "
            "recruiting team reduced average time-to-fill from forty-five to thirty-one days through improved sourcing strategies."
        ),
        "notes": "VP of People will co-present. Highlight DEI progress and remote work policy impact.",
        "has_image": True,
    },
    {  # Slide 13
        "title": "Security and Compliance",
        "body": "Achieved SOC 2 Type II certification. Zero critical security incidents. GDPR audit passed with no findings. ISO 27001 renewal on track.",
        "notes": "CISO prepared a brief demo of the new threat monitoring dashboard if time permits.",
        "has_image": False,
    },
    {  # Slide 14
        "title": "Infrastructure Investments",
        "body": "Migrated 85% of workloads to multi-cloud architecture. Reduced cloud spend by 22% through optimization. New DR site operational.",
        "notes": "",  # NO notes
        "has_image": True,
    },
    {  # Slide 15
        "title": "Customer Support Metrics",
        "body": "Average response time: 2.3 hours. First contact resolution: 78%. CSAT score: 4.6/5.0. Ticket volume down 12% due to self-service improvements.",
        "notes": "New AI chatbot handling 35% of tier-1 inquiries. Full rollout planned for Q4.",
        "has_image": False,
    },
    {  # Slide 16 - overloaded
        "title": "Research and Innovation",
        "body": (
            "The research and innovation division made substantial progress on several strategic initiatives during the quarter. "
            "The machine learning team completed development of our next-generation natural language processing engine, which "
            "demonstrated a forty percent improvement in accuracy compared to the current production model during internal benchmarking. "
            "The advanced analytics group published three peer-reviewed papers at major international conferences, establishing our "
            "thought leadership in predictive modeling for enterprise applications. Patent applications filed increased to twelve, "
            "covering innovations in real-time data processing, anomaly detection, and automated workflow optimization. The prototype "
            "lab delivered two proof-of-concept projects that have been greenlit for product integration in the upcoming release cycle. "
            "Collaboration with three university research programs continued to yield promising results, particularly in the areas of "
            "federated learning and privacy-preserving computation. Budget utilization for R&D was ninety-four percent, reflecting "
            "efficient resource allocation across all active research programs and technology exploration initiatives."
        ),
        "notes": "CTO will demo the NLP engine improvements. Prepare backup slides for technical questions.",
        "has_image": True,
    },
    {  # Slide 17
        "title": "Sustainability Initiatives",
        "body": "Carbon footprint reduced 30%. Office energy consumption down 25%. Launched employee green commute program with 45% participation rate.",
        "notes": "",  # NO notes
        "has_image": True,
    },
    {  # Slide 18
        "title": "Risk Assessment",
        "body": "Key risks: supply chain disruption (medium), regulatory changes (low-medium), talent competition (high), FX exposure (medium).",
        "notes": "Board risk committee reviewed these. Mitigation plans documented in appendix.",
        "has_image": False,
    },
    {  # Slide 19 - long title
        "title": "Strategic Roadmap and Key Initiative Timeline for Fourth Quarter and Full Year Planning Cycle",
        "body": "Q4 priorities: Platform v5.0 launch, APAC market entry, Series C preparation, Enterprise sales acceleration program.",
        "notes": "This is the most critical slide. Allow 10 minutes for discussion.",
        "has_image": True,
    },
    {  # Slide 20
        "title": "Thank You and Next Steps",
        "body": "Next board meeting: December 15, 2025. Action items will be distributed within 48 hours. Contact: strategy@acmetech.com.",
        "notes": "Close with reminder about the holiday party on December 20th.",
        "has_image": True,
    },
]


def create_placeholder_image(path, width=400, height=300, text="Chart", bg_color=(220, 230, 245)):
    """Create a simple placeholder image for slides that need images."""
    from PIL import Image, ImageDraw, ImageFont
    img = Image.new('RGB', (width, height), bg_color)
    draw = ImageDraw.Draw(img)
    # Draw a border
    draw.rectangle([2, 2, width - 3, height - 3], outline=(100, 120, 160), width=2)
    # Draw some placeholder chart bars
    bar_colors = [(70, 130, 180), (60, 179, 113), (218, 165, 32), (205, 92, 92)]
    bar_width = 40
    bar_x_start = 80
    bar_heights = [180, 130, 200, 160]
    for i, (bh, bc) in enumerate(zip(bar_heights, bar_colors)):
        x = bar_x_start + i * 60
        draw.rectangle([x, height - 40 - bh, x + bar_width, height - 40], fill=bc)
    # Add text
    draw.text((width // 2 - 30, 15), text, fill=(60, 60, 80))
    img.save(path)


def create_initial():
    prs = Presentation()

    # Create a temporary chart image for slides that need images
    img_path = f'{WORKDIR}/_chart_placeholder.png'
    create_placeholder_image(img_path)

    for i, sd in enumerate(SLIDE_DATA):
        # Use layout 5 (Blank) for full control
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Add title text box
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8.5), Inches(1.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = sd["title"]
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6D)

        # Add body text box
        body_top = Inches(1.5)
        body_height = Inches(4.0)
        if sd["has_image"]:
            body_width = Inches(5.0)
        else:
            body_width = Inches(8.5)

        body_box = slide.shapes.add_textbox(Inches(0.5), body_top, body_width, body_height)
        tf_body = body_box.text_frame
        tf_body.word_wrap = True
        p_body = tf_body.paragraphs[0]
        p_body.text = sd["body"]
        run_body = p_body.runs[0]
        run_body.font.size = Pt(16)
        run_body.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        # Add image if specified
        if sd["has_image"]:
            slide.shapes.add_picture(
                img_path, Inches(6.0), Inches(1.8), Inches(3.2), Inches(2.4)
            )

        # Add speaker notes if non-empty
        if sd["notes"]:
            slide.notes_slide.notes_text_frame.text = sd["notes"]

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # Clean up temp image
    if os.path.exists(img_path):
        os.remove(img_path)

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
