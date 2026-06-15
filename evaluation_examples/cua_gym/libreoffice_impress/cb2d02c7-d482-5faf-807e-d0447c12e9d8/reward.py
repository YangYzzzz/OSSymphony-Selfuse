"""
Reward Script: Generate a 10-slide presentation with alternating layouts and dark theme
Task ID: impress_gf5_023
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Exactly 10 slides
  Component 2 (0.25): Titles 'Topic 1' through 'Topic 10'
  Component 3 (0.20): Dark background on all slides
  Component 4 (0.30): Alternating layouts (odd=single-column, even=two-column)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_023'


def persist_app_state(domain: str):
    """Best-effort save of any open LibreOffice document."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_slide_bg_color(slide):
    """Return background RGB color of a slide, checking direct fill and master inheritance."""
    try:
        fill = slide.background.fill
        if fill.type == 1:  # SOLID
            return fill.fore_color.rgb
        elif fill.type == 5:  # inherited from master
            try:
                master_fill = slide.slide_layout.slide_master.background.fill
                if master_fill.type == 1:
                    return master_fill.fore_color.rgb
            except Exception:
                pass
    except Exception:
        pass
    return None


def is_dark_color(rgb):
    """Check if an RGB color is 'dark' (luminance < 128)."""
    if rgb is None:
        return False
    r = int(str(rgb)[0:2], 16)
    g = int(str(rgb)[2:4], 16)
    b = int(str(rgb)[4:6], 16)
    luminance = 0.299 * r + 0.587 * g + 0.114 * b
    return luminance < 128


def get_text_shapes(slide):
    """Get all shapes with text frames from a slide."""
    result = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            result.append(shape)
    return result


def get_slide_title_text(slide):
    """Extract the title text from a slide. Looks for a text shape containing 'Topic N'."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            full_text = shape.text_frame.text.strip()
            # The title textbox contains just the topic title
            if full_text.startswith("Topic ") and len(full_text) < 20:
                return full_text
    return None


def count_content_columns(slide):
    """
    Determine if a slide has single-column or two-column content layout.
    Two-column: has 2+ content text boxes side by side (different left positions, similar widths).
    Single-column: has 1 wide content text box.
    Returns 1 or 2.
    """
    content_boxes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            # Skip title textboxes and empty placeholders
            if not text:
                continue
            if text.startswith("Topic ") and len(text) < 20:
                continue
            content_boxes.append(shape)

    if len(content_boxes) >= 2:
        # Check if they are side by side (different left positions)
        lefts = sorted(set(s.left for s in content_boxes))
        if len(lefts) >= 2:
            return 2
    return 1


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Exactly 10 slides (0.25 points)
    try:
        if num_slides == 10:
            print(f"PASS: Component 1 — 10 slides found (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — expected 10 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Titles 'Topic 1' through 'Topic 10' (0.25 points)
    # Each correct title earns 0.025 points
    try:
        title_score = 0.0
        for i in range(min(num_slides, 10)):
            slide = prs.slides[i]
            expected_title = f"Topic {i + 1}"
            actual_title = get_slide_title_text(slide)
            if actual_title == expected_title:
                title_score += 0.025
            else:
                print(f"  Slide {i+1}: expected title '{expected_title}', found '{actual_title}'")

        if title_score >= 0.249:  # all 10 correct (floating point tolerance)
            print(f"PASS: Component 2 — all 10 titles correct (0.25 pts)")
            total_score += 0.25
        elif title_score > 0:
            print(f"PARTIAL: Component 2 — title score {title_score:.3f}/0.25")
            total_score += title_score
        else:
            print(f"FAIL: Component 2 — no correct titles found")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Dark background on all slides (0.20 points)
    # Each slide with dark bg earns 0.02 points
    try:
        dark_count = 0
        for i in range(min(num_slides, 10)):
            slide = prs.slides[i]
            rgb = get_slide_bg_color(slide)
            if is_dark_color(rgb):
                dark_count += 1
            else:
                print(f"  Slide {i+1}: background not dark (color={rgb})")

        bg_score = (dark_count / 10) * 0.20 if num_slides >= 10 else (dark_count / max(num_slides, 1)) * 0.20
        if dark_count == 10:
            print(f"PASS: Component 3 — all 10 slides have dark background (0.20 pts)")
            total_score += 0.20
        elif dark_count > 0:
            print(f"PARTIAL: Component 3 — {dark_count}/10 slides have dark background ({bg_score:.3f} pts)")
            total_score += bg_score
        else:
            print(f"FAIL: Component 3 — no slides have dark background")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Alternating layouts (0.30 points)
    # Odd slides (1,3,5,7,9) should have single-column content
    # Even slides (2,4,6,8,10) should have two-column content
    # Each correct layout earns 0.03 points
    try:
        layout_score = 0.0
        layout_correct = 0
        for i in range(min(num_slides, 10)):
            slide = prs.slides[i]
            cols = count_content_columns(slide)
            slide_num = i + 1
            if slide_num % 2 == 1:
                # Odd slide: expect single-column
                if cols == 1:
                    layout_correct += 1
                    layout_score += 0.03
                else:
                    print(f"  Slide {slide_num} (odd): expected 1 column, found {cols}")
            else:
                # Even slide: expect two-column
                if cols == 2:
                    layout_correct += 1
                    layout_score += 0.03
                else:
                    print(f"  Slide {slide_num} (even): expected 2 columns, found {cols}")

        if layout_correct == 10:
            print(f"PASS: Component 4 — all 10 slides have correct alternating layout (0.30 pts)")
            total_score += 0.30
        elif layout_correct > 0:
            print(f"PARTIAL: Component 4 — {layout_correct}/10 slides correct layout ({layout_score:.3f} pts)")
            total_score += layout_score
        else:
            print(f"FAIL: Component 4 — no slides have correct alternating layout")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/generated_presentation.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
