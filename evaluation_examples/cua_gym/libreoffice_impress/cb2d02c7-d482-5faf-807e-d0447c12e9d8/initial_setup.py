"""
Initial Setup: Open LibreOffice Impress with a blank presentation
Task ID: impress_gf5_023
Domain: libreoffice_impress

Creates a blank presentation with a single empty slide and opens it in Impress.
The agent will then use a Basic macro to generate the 10-slide deck.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_023'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    # Create a blank presentation with one empty slide
    prs = Presentation()
    # Add a single blank slide (layout 6 = Title Only in default template,
    # but we use layout 5 = Blank for a truly empty canvas)
    blank_layout = prs.slide_layouts[6]  # Title Only layout
    slide = prs.slides.add_slide(blank_layout)

    # Leave the slide essentially empty - just a title placeholder with no text
    # This represents the blank state before the macro runs
    if slide.shapes.title:
        slide.shapes.title.text = ""

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress for GUI-ready state
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
