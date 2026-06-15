"""
FINAL REWARD SCRIPT - SUCCESS
Task: I’m cleaning up a huge presentation and noticed slides 10 through 15 still have that old blue backdrop. In LibreOffice Impress, how do I change ONLY those six slides so their background is the built-in “Gray 5%” color and leave the rest of the deck exactly as is?
Generated: 2025-09-10 21:30:53
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
from pptx import Presentation
from pptx.dml.color import RGBColor

"""
Reward Script for LibreOffice Impress Task
Task Recap:
    – Change ONLY slides 10-15 so their background is the built-in “Gray 5%” colour (#F2F2F2)
    – All other slides must remain unchanged (i.e., *not* Gray 5%)
Verification Strategy:
    1. Load the provided .pptx with python-pptx (file existence & loading get NO points – prerequisites).
    2. Examine every slide’s background fill colour.
    3. For slides 10-15 (1-based indices) award credit only when colour == #F2F2F2.
    4. For every other slide award credit only when colour != #F2F2F2 (ensuring they were left alone).

Progressive Scoring:
    • 70% of the score comes from correctly updating the six target slides.
    • 30% comes from correctly *not* altering the remaining slides.
    • Scores are proportional – partial completion yields partial credit.
    • Final score is clipped to [0.0, 1.0] and printed as “REWARD: X.X”.
Anti-Hacking Compliance:
    • No hard-coded True/False; every decision is data-driven.
    • No points for natural conditions like file existence or successful loading.
    • Verification steps can genuinely fail – all checks are falsifiable.
"""

TARGET_FILE = "/home/user/im_cleaning_up_a_huge_presentation_and_noticed_slides_10_through_15_still_have_that_old_blue_backdro_golden.pptx"

# ---------------------------------------------------------------------------
# Helper Functions
# ---------------------------------------------------------------------------

def is_gray5(fore_color):
    """Return True if the given pptx color object represents #F2F2F2 (Gray 5%)."""
    if fore_color is None:
        return False

    # Try direct RGB match first
    try:
        rgb = fore_color.rgb  # May be None if colour comes from a theme
        if rgb and rgb == RGBColor(0xF2, 0xF2, 0xF2):
            return True
    except Exception:
        pass

    # Fallback: sometimes Gray 5% is stored as theme colour with 5% tint – handle common case
    try:
        if hasattr(fore_color, "theme_color") and fore_color.theme_color is not None:
            # In many Office templates Gray 5% maps to theme colour "LT1" with 5% shade
            if str(fore_color.theme_color).upper().endswith("LT1"):
                # Check brightness/tint. Positive = lighter, negative = darker.
                if hasattr(fore_color, "brightness") and abs(fore_color.brightness - 0.05) < 0.0001:
                    return True
    except Exception:
        pass

    return False

# ---------------------------------------------------------------------------
# Core Verification Logic
# ---------------------------------------------------------------------------

def verify_background_changes(file_path: str) -> float:
    print(f"Verifying background changes in: {file_path}")

    # ---------- 0. Preliminary: load file (no points yet) ----------
    if not os.path.exists(file_path):
        print("✗ File not found – cannot verify task")
        print("REWARD: 0.0")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Could not load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    total_slides = len(prs.slides)
    print(f"Total slides detected: {total_slides}")

    if total_slides < 15:
        print("✗ Presentation has fewer than 15 slides – task requirements not satisfied")
        print("REWARD: 0.0")
        return 0.0

    # ---------- 1. Analyse each slide background ----------
    target_indices = list(range(10, 16))  # Slides 10-15 inclusive (1-based)

    target_correct = 0          # How many of 10-15 have correct colour
    normal_correct = 0          # How many non-targets kept original colour (i.e., not Gray 5%)
    normal_total = total_slides - len(target_indices)

    for idx, slide in enumerate(prs.slides, start=1):
        # Obtain fill and colour details; pptx background always has .fill
        fill = slide.background.fill
        fore_color = None
        if fill is not None:
            try:
                fore_color = fill.fore_color
            except Exception:
                pass

        is_g5 = is_gray5(fore_color)

        if idx in target_indices:
            if is_g5:
                target_correct += 1
                print(f"✓ Slide {idx}: Gray 5% correctly applied")
            else:
                print(f"✗ Slide {idx}: Background is NOT Gray 5% as required")
        else:
            if not is_g5:
                normal_correct += 1
            else:
                print(f"✗ Slide {idx}: Should remain unchanged but is Gray 5% – incorrectly altered")

    # ---------- 2. Scoring ----------
    target_ratio = target_correct / len(target_indices)
    normal_ratio = normal_correct / normal_total if normal_total else 1.0

    score = target_ratio * 0.7 + normal_ratio * 0.3  # Weighted composite score
    score = round(min(max(score, 0.0), 1.0), 4)      # Clamp & round for neatness

    print(f"Changed slides correct: {target_correct}/{len(target_indices)} → {target_ratio:.0%}")
    print(f"Unchanged slides correct: {normal_correct}/{normal_total} → {normal_ratio:.0%}")
    print(f"FINAL SCORE: {score}")
    print(f"REWARD: {score}")

    return score

# ---------------------------------------------------------------------------
# Execute verification when run as a script
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    verify_background_changes(TARGET_FILE)

