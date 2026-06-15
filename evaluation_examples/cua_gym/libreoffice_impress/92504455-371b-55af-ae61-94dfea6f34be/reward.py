"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 5, I want to add some vertical text to stand out more. How do I insert a text box with the heading 'Q2 RESULTS' that goes from top to bottom instead of left to right?
Generated: 2025-08-07 13:14:40
Status: success
Model: o4-mini
Total Steps: 4
"""

import os
from pptx import Presentation

def verify_task(file_path):
    print("Starting task verification for vertical text insertion...")
    total_score = 0.0
    max_score = 1.0

    # 1. Check file existence and load
    print("\n1. Checking file existence and load")
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print(f"REWARD: {total_score}")
        return total_score
    print(f"✓ File exists: {file_path} (0.2)")
    total_score += 0.2

    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        print(f"REWARD: {total_score}")
        return total_score

    # 2. Check slide 5 exists
    print("\n2. Verifying slide count for slide 5 existence")
    if len(prs.slides) >= 5:
        print("✓ Slide 5 exists (0.2)")
        total_score += 0.2
    else:
        print(f"✗ Slide count is {len(prs.slides)}, slide 5 missing")

    # 3. Verify presence of text box with 'Q2 RESULTS'
    print("\n3. Verifying 'Q2 RESULTS' text box presence")
    target_shape = None
    slide5 = prs.slides[4] if len(prs.slides) >= 5 else None
    if slide5:
        for shape in slide5.shapes:
            if hasattr(shape, 'text') and shape.text.strip().upper() == 'Q2 RESULTS':
                target_shape = shape
                break
    if target_shape:
        print("✓ Found text box 'Q2 RESULTS' (0.3)")
        total_score += 0.3
    else:
        print("✗ 'Q2 RESULTS' text box not found on slide 5")

    # 4. Verify vertical orientation via rotation attribute
    print("\n4. Verifying vertical orientation (rotation) of the text box")
    if target_shape:
        try:
            ns = {
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
            }
            xfrm = target_shape._element.find('.//a:xfrm', ns)
            if xfrm is not None:
                rot_val = xfrm.get('rot')
                if rot_val is not None:
                    try:
                        rot_int = int(rot_val)
                    except ValueError:
                        rot_int = 0
                    if rot_int != 0:
                        print(f"✓ Rotation attribute found: rot={rot_int} (0.3)")
                        total_score += 0.3
                    else:
                        print(f"✗ Rotation is zero, text is not vertical: rot={rot_int}")
                else:
                    print("✗ No rotation attribute on xfrm element")
            else:
                print("✗ No xfrm element found for shape")
        except Exception as e:
            print(f"✗ Error verifying rotation: {e}")
    else:
        print("✗ Skipping rotation check because 'Q2 RESULTS' shape not found")

    # Final scoring
    final_score = min(total_score, max_score)
    print(f"\nScore summary: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score

if __name__ == '__main__':
    file_path = '/home/user/on_slide_5_i_want_to_add_some_vertical_text_to_stand_out_more_how_do_i_insert_a_text_box_with_the_he.pptx'
    verify_task(file_path)
