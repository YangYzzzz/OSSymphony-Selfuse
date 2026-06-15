"""
Initial Setup: Visual Storytelling Deck - Titles at Top (pre-task state)
Task ID: osworld_impress_title_position_bottom_009
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_position_bottom_009'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard widescreen: 13.33 x 7.5 inches
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_w = prs.slide_width   # EMU
    slide_h = prs.slide_height  # EMU

    # --- Slide data: 6 visual storytelling slides ---
    slides_data = [
        {
            "title": "The Journey Begins",
            "subtitle": "A story of discovery and transformation across distant lands",
            "bg_color": RGBColor(0x1A, 0x1A, 0x2E),
            "title_color": RGBColor(0xFF, 0xFF, 0xFF),
            "sub_color": RGBColor(0xCC, 0xCC, 0xCC),
        },
        {
            "title": "Crossing the Mountains",
            "subtitle": "Navigating the rugged terrain of the Northern Highlands",
            "bg_color": RGBColor(0x16, 0x21, 0x3E),
            "title_color": RGBColor(0xF5, 0xF5, 0xF5),
            "sub_color": RGBColor(0xBB, 0xBB, 0xCC),
        },
        {
            "title": "The Hidden Valley",
            "subtitle": "Where ancient rivers carve stories into the stone",
            "bg_color": RGBColor(0x0F, 0x3D, 0x2E),
            "title_color": RGBColor(0xEE, 0xFF, 0xEE),
            "sub_color": RGBColor(0xAA, 0xDD, 0xBB),
        },
        {
            "title": "City of Light",
            "subtitle": "The metropolitan heartbeat of modern civilization",
            "bg_color": RGBColor(0x2C, 0x22, 0x10),
            "title_color": RGBColor(0xFF, 0xF0, 0xCC),
            "sub_color": RGBColor(0xDD, 0xCC, 0xAA),
        },
        {
            "title": "Return to the Sea",
            "subtitle": "The endless horizon where sky meets water in silence",
            "bg_color": RGBColor(0x05, 0x28, 0x4A),
            "title_color": RGBColor(0xDD, 0xF0, 0xFF),
            "sub_color": RGBColor(0xAA, 0xCC, 0xEE),
        },
        {
            "title": "The Story Continues",
            "subtitle": "Every ending holds the seed of a new beginning",
            "bg_color": RGBColor(0x2E, 0x1A, 0x35),
            "title_color": RGBColor(0xF5, 0xEE, 0xFF),
            "sub_color": RGBColor(0xCC, 0xBB, 0xDD),
        },
    ]

    for slide_info in slides_data:
        # Use blank layout for full control
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

        # --- Full-bleed background color (simulates full-bleed image) ---
        bg_fill = slide.background.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = slide_info["bg_color"]

        # --- Title textbox at TOP of slide (within top half) ---
        # Top: ~5% from top, height ~15% of slide
        title_top = Emu(int(slide_h * 0.05))
        title_left = Emu(int(slide_w * 0.05))
        title_width = Emu(int(slide_w * 0.90))
        title_height = Emu(int(slide_h * 0.15))

        title_box = slide.shapes.add_textbox(
            title_left, title_top, title_width, title_height
        )
        tf = title_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = slide_info["title"]
        run.font.name = "Arial"
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = slide_info["title_color"]

        # --- Subtitle/caption textbox (below title, still upper area) ---
        sub_top = Emu(int(slide_h * 0.22))
        sub_left = Emu(int(slide_w * 0.05))
        sub_width = Emu(int(slide_w * 0.70))
        sub_height = Emu(int(slide_h * 0.10))

        sub_box = slide.shapes.add_textbox(
            sub_left, sub_top, sub_width, sub_height
        )
        tf2 = sub_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        run2 = p2.add_run()
        run2.text = slide_info["subtitle"]
        run2.font.name = "Arial"
        run2.font.size = Pt(18)
        run2.font.bold = False
        run2.font.color.rgb = slide_info["sub_color"]

        # --- Small decorative accent bar near top ---
        accent_top = Emu(int(slide_h * 0.02))
        accent_left = Emu(int(slide_w * 0.05))
        accent_width = Emu(int(slide_w * 0.08))
        accent_height = Emu(int(slide_h * 0.008))

        accent_box = slide.shapes.add_textbox(
            accent_left, accent_top, accent_width, accent_height
        )
        accent_fill = accent_box.fill
        accent_fill.solid()
        accent_fill.fore_color.rgb = RGBColor(0xFF, 0xCC, 0x00)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
