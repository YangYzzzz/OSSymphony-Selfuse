"""
Reward Script: Move the title textbox on all 6 slides to the bottom of each slide
Task ID: osworld_impress_title_position_bottom_009
Domain: libreoffice_impress
Scoring:
  - Component 1: At least 1 title (TextBox 1) is in the bottom half of its slide (0.3 pts)
  - Component 2: At least 3 titles (TextBox 1) are in the bottom half of their slides (0.3 pts)
  - Component 3: All 6 titles (TextBox 1) are in the bottom half of their slides (0.4 pts)
Total: 1.0

Verification logic:
  The task requires moving ALL title textboxes (named 'TextBox 1', containing the main
  slide titles) to the BOTTOM half of each slide.

  Slide height = 6858000 EMU (7.5 inches).
  Bottom half means: shape.top >= slide_height / 2 = 3429000 EMU (3.75 inches).

  Initial state: all TextBox 1 shapes have top=342900 (0.38 inches) -> in top half
  Golden state: all TextBox 1 shapes have top=5143500 (5.62 inches) -> in bottom half
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_position_bottom_009'

# Title shape names to look for on each slide
TITLE_SHAPE_NAMES = {'TextBox 1', 'Title 1', 'Title'}

# Title text strings from the 6 slides
EXPECTED_TITLES = {
    'The Journey Begins',
    'Crossing the Mountains',
    'The Hidden Valley',
    'City of Light',
    'Return to the Sea',
    'The Story Continues',
}


def is_title_shape(shape):
    """
    Determine if a shape is the main title shape.
    Prioritize 'TextBox 1' by name, fall back to checking text content.
    """
    if not shape.has_text_frame:
        return False
    if shape.name in TITLE_SHAPE_NAMES:
        return True
    return False


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Checks that title textboxes on all 6 slides have been moved to the bottom half.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: verify we have exactly 6 slides
    num_slides = len(prs.slides)
    if num_slides != 6:
        print(f"PRECONDITION FAIL: Expected 6 slides, found {num_slides}. Cannot score.")
        print("REWARD: 0.0")
        return 0.0

    slide_height = prs.slide_height
    half_height = slide_height / 2

    print(f"Slide height: {slide_height} EMU ({slide_height/914400:.2f} inches)")
    print(f"Bottom half threshold (top must be >= ): {half_height} EMU ({half_height/914400:.2f} inches)")
    print()

    # Collect info about title shapes across all 6 slides
    slides_with_title_in_bottom = 0
    total_slides_with_title = 0
    slide_details = []

    for i, slide in enumerate(prs.slides):
        title_shape = None
        for shape in slide.shapes:
            if is_title_shape(shape):
                # Check if it's a likely title (TextBox 1 or has known title text)
                text = shape.text_frame.text.strip()
                if shape.name == 'TextBox 1':
                    title_shape = shape
                    break

        if title_shape is not None:
            total_slides_with_title += 1
            top = title_shape.top
            in_bottom = (top >= half_height)
            slide_details.append({
                'slide_num': i + 1,
                'shape_name': title_shape.name,
                'text': title_shape.text_frame.text[:40],
                'top': top,
                'top_inches': top / 914400,
                'in_bottom_half': in_bottom,
            })
            if in_bottom:
                slides_with_title_in_bottom += 1
            print(f"Slide {i+1}: '{title_shape.text_frame.text[:30]}' top={top} ({top/914400:.2f}in) -> {'BOTTOM' if in_bottom else 'TOP'} half")
        else:
            print(f"Slide {i+1}: No TextBox 1 title shape found")
            slide_details.append({
                'slide_num': i + 1,
                'shape_name': None,
                'in_bottom_half': False,
            })

    print()
    print(f"Titles in bottom half: {slides_with_title_in_bottom} / {total_slides_with_title} (out of 6 slides)")
    print()

    # Component 1: At least 1 title moved to bottom half (0.3 pts)
    # This FAILS on initial_env (all titles at top=342900, top half) and
    # PASSES on golden_env (all titles at top=5143500, bottom half)
    try:
        if slides_with_title_in_bottom >= 1:
            print(f"PASS: Component 1 — At least 1 title in bottom half ({slides_with_title_in_bottom} found) (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 — No titles found in bottom half (expected >= 1)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Majority of titles (3+) moved to bottom half (0.3 pts)
    # Additional partial credit — only awarded if more than half of titles are moved
    try:
        if slides_with_title_in_bottom >= 3:
            print(f"PASS: Component 2 — At least 3 titles in bottom half ({slides_with_title_in_bottom} found) (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 — Only {slides_with_title_in_bottom}/6 titles in bottom half (expected >= 3)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: All 6 titles moved to bottom half (0.4 pts)
    # Full credit only awarded when ALL 6 titles are in the bottom half
    try:
        if slides_with_title_in_bottom == 6:
            print(f"PASS: Component 3 — All 6 titles are in the bottom half (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 3 — Only {slides_with_title_in_bottom}/6 titles in bottom half (expected all 6)")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Run against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
