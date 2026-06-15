"""
Reward Script: Executive Summary Slides (Market Opportunity, Business Model, Financial Projections)
Task ID: impress_sales_059
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 2 has a doughnut chart (0.20)
  Component 2: Slide 2 doughnut chart has 3 labeled segments for TAM/SAM/SOM (0.15)
  Component 3: Slide 3 has flow diagram shapes with correct labels (0.20)
  Component 4: Slide 3 has connectors linking shapes (0.10)
  Component 5: Slide 4 has a chart present (0.15)
  Component 6: Slide 4 chart is a combo chart with bar + line plots (0.20)
"""

import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_059'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]  # Slide 2: Market Opportunity
    slide3 = prs.slides[2]  # Slide 3: Business Model
    slide4 = prs.slides[3]  # Slide 4: Financial Projections

    # =========================================================================
    # Component 1: Slide 2 has a doughnut chart (0.20 points)
    # Initial: no chart; Golden: doughnut chart present
    # =========================================================================
    try:
        doughnut_chart = None
        for shape in slide2.shapes:
            if hasattr(shape, 'has_chart') and shape.has_chart:
                chart = shape.chart
                # DOUGHNUT enum is -4120
                if chart.chart_type == -4120:
                    doughnut_chart = chart
                    break

        if doughnut_chart is not None:
            print(f"PASS: Component 1 -- Slide 2 has a doughnut chart (0.20 pts)")
            total_score += 0.20
        else:
            # Check if there's any chart at all on slide 2
            any_chart = any(hasattr(s, 'has_chart') and s.has_chart for s in slide2.shapes)
            if any_chart:
                print(f"FAIL: Component 1 -- Slide 2 has a chart but it is not a doughnut chart")
            else:
                print(f"FAIL: Component 1 -- Slide 2 has no chart")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # =========================================================================
    # Component 2: Doughnut chart has 3 labeled segments for TAM/SAM/SOM (0.15 points)
    # Initial: no chart; Golden: categories contain TAM, SAM, SOM labels
    # =========================================================================
    try:
        if doughnut_chart is not None:
            plot = doughnut_chart.plots[0]
            categories = [str(c) for c in plot.categories]
            num_series_values = list(plot.series[0].values) if len(plot.series) > 0 else []

            # Check that there are 3 categories mentioning TAM, SAM, SOM
            has_tam = any('TAM' in c.upper() or 'TAM' in c for c in categories)
            has_sam = any('SAM' in c.upper() or 'SAM' in c for c in categories)
            has_som = any('SOM' in c.upper() or 'SOM' in c for c in categories)
            has_three = len(categories) >= 3

            if has_tam and has_sam and has_som and has_three:
                print(f"PASS: Component 2 -- Doughnut chart has TAM/SAM/SOM categories: {categories} (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 2 -- Expected TAM/SAM/SOM categories, found: {categories}")
        else:
            print(f"FAIL: Component 2 -- No doughnut chart to check categories")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # =========================================================================
    # Component 3: Slide 3 has flow diagram shapes with correct labels (0.20 points)
    # Initial: no shapes beyond title; Golden: rounded rectangles with
    # Revenue Streams, Subscriptions, Professional Services, Marketplace
    # =========================================================================
    try:
        auto_shape_texts = []
        for shape in slide3.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if shape.has_text_frame:
                    # Normalize vertical tab and newlines
                    text = shape.text_frame.text.replace('\x0b', ' ').replace('\n', ' ').strip()
                    auto_shape_texts.append(text.lower())

        # Check for key labels
        required_labels = ['subscriptions', 'professional', 'marketplace']
        found_labels = []
        for label in required_labels:
            if any(label in t for t in auto_shape_texts):
                found_labels.append(label)

        # Also check for "revenue streams" as the parent node
        has_revenue_streams = any('revenue' in t for t in auto_shape_texts)

        if len(found_labels) >= 3 and has_revenue_streams:
            print(f"PASS: Component 3 -- Slide 3 has flow diagram with all required labels: {auto_shape_texts} (0.20 pts)")
            total_score += 0.20
        elif len(found_labels) >= 2:
            print(f"PARTIAL: Component 3 -- Found {len(found_labels)}/3 required labels + revenue_streams={has_revenue_streams}: {auto_shape_texts} (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 3 -- Expected flow diagram shapes, found auto_shapes: {auto_shape_texts}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # =========================================================================
    # Component 4: Slide 3 has connectors linking shapes (0.10 points)
    # Initial: no connectors; Golden: >= 2 connector/line shapes
    # =========================================================================
    try:
        connector_count = 0
        for shape in slide3.shapes:
            # LINE (9) or FREEFORM connectors
            if shape.shape_type == MSO_SHAPE_TYPE.LINE or shape.shape_type == 9:
                connector_count += 1

        if connector_count >= 2:
            print(f"PASS: Component 4 -- Slide 3 has {connector_count} connectors (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 4 -- Expected >= 2 connectors on slide 3, found {connector_count}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # =========================================================================
    # Component 5: Slide 4 has a chart present (0.15 points)
    # Initial: no chart; Golden: chart present
    # =========================================================================
    try:
        slide4_chart = None
        for shape in slide4.shapes:
            if hasattr(shape, 'has_chart') and shape.has_chart:
                slide4_chart = shape.chart
                break

        if slide4_chart is not None:
            print(f"PASS: Component 5 -- Slide 4 has a chart (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 5 -- Slide 4 has no chart")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # =========================================================================
    # Component 6: Slide 4 chart is a combo chart with bar + line plots (0.20 points)
    # Initial: no chart; Golden: 2 plots (BarPlot + LinePlot)
    # =========================================================================
    try:
        if slide4_chart is not None:
            num_plots = len(slide4_chart.plots)
            plot_types = [type(p).__name__ for p in slide4_chart.plots]

            has_bar = any('Bar' in pt for pt in plot_types)
            has_line = any('Line' in pt for pt in plot_types)

            if has_bar and has_line and num_plots >= 2:
                print(f"PASS: Component 6 -- Slide 4 has combo chart (bar+line): {plot_types} (0.20 pts)")
                total_score += 0.20
            elif has_bar or has_line:
                print(f"PARTIAL: Component 6 -- Slide 4 chart has only {plot_types}, expected both bar+line (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 6 -- Slide 4 chart plots are {plot_types}, expected bar+line combo")
        else:
            print(f"FAIL: Component 6 -- No chart on slide 4 to check combo type")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
