"""
Initial Setup: Create a 10-slide executive pitch deck with title-only slides 2-4
Task ID: impress_sales_059
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_059'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a styled text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    # Standard 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== Slide 1: Title Slide =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "NexGen Analytics"
    slide1.placeholders[1].text = "Executive Summary & Investment Pitch\nQ2 2025 Series B"

    # ===== Slide 2: Market Opportunity (TITLE ONLY - no content) =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    add_textbox(slide2, Inches(0.5), Inches(0.3), Inches(10), Inches(1),
                "Market Opportunity", font_size=36, bold=True,
                color=RGBColor(0x1F, 0x3A, 0x5F))

    # ===== Slide 3: Business Model (TITLE ONLY - no content) =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    add_textbox(slide3, Inches(0.5), Inches(0.3), Inches(10), Inches(1),
                "Business Model", font_size=36, bold=True,
                color=RGBColor(0x1F, 0x3A, 0x5F))

    # ===== Slide 4: Financial Projections (TITLE ONLY - no content) =====
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    add_textbox(slide4, Inches(0.5), Inches(0.3), Inches(10), Inches(1),
                "Financial Projections", font_size=36, bold=True,
                color=RGBColor(0x1F, 0x3A, 0x5F))

    # ===== Slide 5: Product Overview =====
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide5, Inches(0.5), Inches(0.3), Inches(10), Inches(1),
                "Product Overview", font_size=36, bold=True,
                color=RGBColor(0x1F, 0x3A, 0x5F))
    add_textbox(slide5, Inches(0.7), Inches(1.5), Inches(11), Inches(1),
                "AI-Powered Business Intelligence Platform", font_size=24,
                color=RGBColor(0x33, 0x33, 0x33))
    features = [
        "Real-time data analytics with predictive modeling",
        "Natural language query interface for non-technical users",
        "Custom dashboard builder with 50+ visualization templates",
        "Automated anomaly detection and alerting system",
        "Enterprise-grade security with SOC 2 Type II compliance",
    ]
    for i, feat in enumerate(features):
        add_textbox(slide5, Inches(1.0), Inches(2.5 + i * 0.7), Inches(10), Inches(0.6),
                    f"• {feat}", font_size=16, color=RGBColor(0x44, 0x44, 0x44))

    # ===== Slide 6: Competitive Landscape =====
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide6, Inches(0.5), Inches(0.3), Inches(10), Inches(1),
                "Competitive Landscape", font_size=36, bold=True,
                color=RGBColor(0x1F, 0x3A, 0x5F))
    competitors = [
        ("Tableau / Power BI", "Legacy BI tools — no predictive capabilities"),
        ("Looker / Mode", "SQL-first — steep learning curve for business users"),
        ("ThoughtSpot", "Search-based BI — limited customization options"),
        ("NexGen Analytics", "AI-native platform — predictive + accessible + customizable"),
    ]
    for i, (name, desc) in enumerate(competitors):
        add_textbox(slide6, Inches(0.7), Inches(1.8 + i * 1.0), Inches(4), Inches(0.5),
                    name, font_size=18, bold=True,
                    color=RGBColor(0x1F, 0x3A, 0x5F) if i == 3 else RGBColor(0x55, 0x55, 0x55))
        add_textbox(slide6, Inches(5.5), Inches(1.8 + i * 1.0), Inches(7), Inches(0.5),
                    desc, font_size=16, color=RGBColor(0x44, 0x44, 0x44))

    # ===== Slide 7: Go-to-Market Strategy =====
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide7, Inches(0.5), Inches(0.3), Inches(10), Inches(1),
                "Go-to-Market Strategy", font_size=36, bold=True,
                color=RGBColor(0x1F, 0x3A, 0x5F))
    strategies = [
        "Phase 1: Enterprise direct sales targeting Fortune 500 ($250K+ ACV)",
        "Phase 2: Mid-market expansion via channel partnerships (Deloitte, Accenture)",
        "Phase 3: Self-serve SMB tier launch with PLG motion ($49/user/month)",
        "Phase 4: International expansion — EMEA (London HQ), APAC (Singapore HQ)",
    ]
    for i, strat in enumerate(strategies):
        add_textbox(slide7, Inches(0.7), Inches(1.8 + i * 1.1), Inches(11), Inches(0.8),
                    strat, font_size=18, color=RGBColor(0x33, 0x33, 0x33))

    # ===== Slide 8: Team =====
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide8, Inches(0.5), Inches(0.3), Inches(10), Inches(1),
                "Leadership Team", font_size=36, bold=True,
                color=RGBColor(0x1F, 0x3A, 0x5F))
    team = [
        ("Dr. Elena Vasquez", "CEO & Co-founder", "Ex-Google AI, Stanford PhD, 15yr experience"),
        ("Raj Patel", "CTO & Co-founder", "Ex-Meta Engineering Director, MIT MS CS"),
        ("Sarah Kim", "VP Product", "Ex-Salesforce, 10yr enterprise SaaS product leadership"),
        ("James O'Brien", "VP Sales", "Ex-Snowflake, built $200M ARR sales org"),
        ("Priya Sharma", "CFO", "Ex-Stripe Finance, Goldman Sachs alum, CPA"),
    ]
    for i, (name, title, bio) in enumerate(team):
        add_textbox(slide8, Inches(0.7), Inches(1.5 + i * 1.0), Inches(3.5), Inches(0.5),
                    name, font_size=18, bold=True, color=RGBColor(0x1F, 0x3A, 0x5F))
        add_textbox(slide8, Inches(4.5), Inches(1.5 + i * 1.0), Inches(2.5), Inches(0.5),
                    title, font_size=16, bold=True, color=RGBColor(0x55, 0x55, 0x55))
        add_textbox(slide8, Inches(7.5), Inches(1.5 + i * 1.0), Inches(5.5), Inches(0.5),
                    bio, font_size=14, color=RGBColor(0x66, 0x66, 0x66))

    # ===== Slide 9: Traction & Milestones =====
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide9, Inches(0.5), Inches(0.3), Inches(10), Inches(1),
                "Traction & Milestones", font_size=36, bold=True,
                color=RGBColor(0x1F, 0x3A, 0x5F))
    milestones = [
        "Q1 2024: Product launch — 12 beta customers onboarded",
        "Q3 2024: $1.2M ARR reached — 45 paying enterprise accounts",
        "Q4 2024: Series A closed ($18M @ $90M pre-money, led by Sequoia)",
        "Q1 2025: 120 enterprise customers, $4.8M ARR, 95% gross retention",
        "Q2 2025: Partnership with AWS Marketplace — 3x inbound pipeline",
    ]
    for i, ms in enumerate(milestones):
        add_textbox(slide9, Inches(0.7), Inches(1.8 + i * 0.9), Inches(11), Inches(0.7),
                    ms, font_size=18, color=RGBColor(0x33, 0x33, 0x33))

    # ===== Slide 10: The Ask =====
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide10, Inches(0.5), Inches(0.3), Inches(10), Inches(1),
                "The Ask", font_size=36, bold=True,
                color=RGBColor(0x1F, 0x3A, 0x5F))
    add_textbox(slide10, Inches(1.0), Inches(1.8), Inches(10), Inches(1),
                "Series B: $50M at $350M Pre-Money Valuation",
                font_size=28, bold=True, color=RGBColor(0x2E, 0x75, 0xB6))
    uses = [
        "Engineering & Product: $20M — Scale platform, hire 40 engineers",
        "Sales & Marketing: $15M — Expand GTM team, launch EMEA operations",
        "Infrastructure: $10M — Multi-cloud deployment, data center expansion",
        "Working Capital & G&A: $5M — Operational runway and legal/compliance",
    ]
    for i, use in enumerate(uses):
        add_textbox(slide10, Inches(1.0), Inches(3.2 + i * 0.8), Inches(10), Inches(0.6),
                    use, font_size=16, color=RGBColor(0x44, 0x44, 0x44))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
