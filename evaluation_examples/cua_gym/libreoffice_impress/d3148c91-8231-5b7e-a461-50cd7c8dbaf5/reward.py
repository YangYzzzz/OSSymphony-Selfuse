"""
Reward Script: Tool Comparison Presentation (Asana vs Trello vs Monday)
Task ID: impress_wf_033
Domain: libreoffice_impress
Scoring:
  C1 (0.15) - File exists on Desktop with exactly 8 slides
  C2 (0.10) - Slide 1 title contains 'Asana vs Trello vs Monday'
  C3 (0.10) - Slide 2 lists 6 evaluation criteria
  C4 (0.20) - Slides 3-5 have pros in green and cons in red text
  C5 (0.20) - Slide 6 has comparison table with 3 tools, 6 criteria, checkmarks/X
  C6 (0.10) - Slide 7 has a grouped bar chart
  C7 (0.15) - Background #F5F5F5 and teal #009688 accent on slides
"""

import os

WORKDIR = '/home/user'
FILE_PATH = os.path.join(WORKDIR, 'Desktop', 'Tool_Comparison.pptx')


def get_all_text_from_slide(slide):
    """Recursively get all text from a slide, including grouped shapes."""
    texts = []
    def extract(shape):
        if hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        try:
                            rgb = str(run.font.color.rgb) if run.font.color.type is not None else None
                        except:
                            rgb = None
                        texts.append((run.text.strip(), rgb, run.font.bold))
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                extract(sub)
    for shape in slide.shapes:
        extract(shape)
    return texts


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: File has exactly 8 slides (0.15 points)
    try:
        num_slides = len(prs.slides)
        if num_slides == 8:
            print(f"PASS: Component 1 — 8 slides found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — expected 8 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 1 title contains 'Asana vs Trello vs Monday' (0.10 points)
    try:
        slide1 = prs.slides[0]
        slide1_texts = get_all_text_from_slide(slide1)
        all_text = " ".join([t[0] for t in slide1_texts]).lower()
        if "asana" in all_text and "trello" in all_text and "monday" in all_text:
            print(f"PASS: Component 2 — Slide 1 title references all 3 tools (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 — Slide 1 text missing tool names. Found: {all_text[:100]}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 2 has 6 evaluation criteria listed (0.10 points)
    try:
        slide2 = prs.slides[1]
        slide2_texts = get_all_text_from_slide(slide2)
        # Count distinct criteria-related text entries (non-title, non-empty)
        criteria_keywords = ["ease", "task", "collaborat", "report", "integrat", "pric"]
        criteria_found = 0
        all_slide2_text = " ".join([t[0] for t in slide2_texts]).lower()
        for kw in criteria_keywords:
            if kw in all_slide2_text:
                criteria_found += 1
        if criteria_found >= 6:
            print(f"PASS: Component 3 — All 6 evaluation criteria found on Slide 2 (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 3 — Only {criteria_found}/6 criteria keywords found on Slide 2")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slides 3-5 have pros in green and cons in red text (0.20 points)
    try:
        green_variants = {"2E7D32", "00FF00", "008000", "009688", "4CAF50", "388E3C", "1B5E20", "66BB6A", "009600"}
        red_variants = {"C62828", "FF0000", "CC0000", "D32F2F", "B71C1C", "E53935", "F44336", "FF1744", "DD0000"}

        def is_green(rgb_str):
            if rgb_str is None:
                return False
            rgb_upper = rgb_str.upper()
            if rgb_upper in green_variants:
                return True
            # Also check if green channel dominates
            try:
                r = int(rgb_upper[0:2], 16)
                g = int(rgb_upper[2:4], 16)
                b = int(rgb_upper[4:6], 16)
                if g > r and g > b and g >= 100:
                    return True
            except:
                pass
            return False

        def is_red(rgb_str):
            if rgb_str is None:
                return False
            rgb_upper = rgb_str.upper()
            if rgb_upper in red_variants:
                return True
            try:
                r = int(rgb_upper[0:2], 16)
                g = int(rgb_upper[2:4], 16)
                b = int(rgb_upper[4:6], 16)
                if r > g and r > b and r >= 150:
                    return True
            except:
                pass
            return False

        slides_with_green_pros = 0
        slides_with_red_cons = 0
        for idx in [2, 3, 4]:  # Slides 3, 4, 5
            if idx >= len(prs.slides):
                continue
            slide = prs.slides[idx]
            texts = get_all_text_from_slide(slide)
            has_green = False
            has_red = False
            for text, rgb, bold in texts:
                if is_green(rgb):
                    has_green = True
                if is_red(rgb):
                    has_red = True
            if has_green:
                slides_with_green_pros += 1
            if has_red:
                slides_with_red_cons += 1

        if slides_with_green_pros >= 3 and slides_with_red_cons >= 3:
            print(f"PASS: Component 4 — All 3 tool slides have green pros and red cons (0.20 pts)")
            total_score += 0.20
        elif slides_with_green_pros >= 2 and slides_with_red_cons >= 2:
            print(f"PARTIAL: Component 4 — {slides_with_green_pros}/3 green, {slides_with_red_cons}/3 red (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 4 — green pros: {slides_with_green_pros}/3, red cons: {slides_with_red_cons}/3")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slide 6 has comparison table with 3 tools and 6 criteria with checkmarks/X (0.20 points)
    try:
        slide6 = prs.slides[5]
        table_found = False
        table_ok = False
        for shape in slide6.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_found = True
                table = shape.table
                num_rows = len(table.rows)
                num_cols = len(table.columns)
                print(f"  Table found: {num_rows} rows x {num_cols} cols")

                # Need at least 4 columns (criteria + 3 tools) and 7 rows (header + 6 criteria)
                has_tools = num_cols >= 4
                has_criteria = num_rows >= 7

                # Check for checkmarks and X symbols in table content
                checkmark_count = 0
                x_count = 0
                for r in range(1, num_rows):  # skip header
                    for c in range(1, num_cols):  # skip criteria column
                        cell_text = table.cell(r, c).text.strip()
                        if any(ch in cell_text for ch in ["\u2713", "\u2714", "\u2705", "\u2611"]):
                            checkmark_count += 1
                        if any(ch in cell_text for ch in ["\u2717", "\u2718", "\u2716", "\u274C", "\u2715"]):
                            x_count += 1

                has_symbols = checkmark_count >= 3 and x_count >= 2

                if has_tools and has_criteria and has_symbols:
                    print(f"PASS: Component 5 — Table with {num_cols-1} tools, {num_rows-1} criteria, {checkmark_count} checkmarks, {x_count} X marks (0.20 pts)")
                    total_score += 0.20
                    table_ok = True
                elif has_tools and has_criteria:
                    print(f"PARTIAL: Component 5 — Table structure OK but symbols incomplete: {checkmark_count} checkmarks, {x_count} X marks (0.10 pts)")
                    total_score += 0.10
                    table_ok = True
                else:
                    print(f"FAIL: Component 5 — Table structure: tools={has_tools}, criteria={has_criteria}")
                break

        if not table_found:
            print(f"FAIL: Component 5 — No table found on Slide 6")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Slide 7 has a grouped bar chart (0.10 points)
    try:
        slide7 = prs.slides[6]
        chart_found = False
        for shape in slide7.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                chart = shape.chart
                chart_type = chart.chart_type
                series_count = len(chart.series)
                print(f"  Chart found: type={chart_type}, series={series_count}")

                # BAR_CLUSTERED = 57, COLUMN_CLUSTERED = 51
                # A "grouped bar chart" could be either bar or column clustered
                from pptx.enum.chart import XL_CHART_TYPE
                is_grouped = chart_type in (
                    XL_CHART_TYPE.BAR_CLUSTERED,
                    XL_CHART_TYPE.COLUMN_CLUSTERED,
                )
                has_multiple_series = series_count >= 2

                if is_grouped and has_multiple_series:
                    print(f"PASS: Component 6 — Grouped bar chart with {series_count} series (0.10 pts)")
                    total_score += 0.10
                    chart_found = True
                elif chart_type is not None:
                    # Any chart is partial credit
                    print(f"PARTIAL: Component 6 — Chart found but type={chart_type}, series={series_count} (0.05 pts)")
                    total_score += 0.05
                    chart_found = True
                break

        if not chart_found:
            print(f"FAIL: Component 6 — No chart found on Slide 7")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Background #F5F5F5 and teal #009688 accents (0.15 points)
    try:
        bg_ok_count = 0
        teal_found = False

        # Check background color on multiple slides
        for idx in range(min(len(prs.slides), 8)):
            slide = prs.slides[idx]
            fill = slide.background.fill
            if fill.type == 1:  # SOLID
                bg_rgb = str(fill.fore_color.rgb).upper()
                if bg_rgb == "F5F5F5":
                    bg_ok_count += 1

        # Check for teal accent (#009688) in text across slides
        for idx in range(min(len(prs.slides), 8)):
            slide = prs.slides[idx]
            texts = get_all_text_from_slide(slide)
            for text, rgb, bold in texts:
                if rgb is not None and rgb.upper() == "009688":
                    teal_found = True
                    break
            if teal_found:
                break

        if bg_ok_count >= 6 and teal_found:
            print(f"PASS: Component 7 — {bg_ok_count}/8 slides with #F5F5F5 bg, teal accent found (0.15 pts)")
            total_score += 0.15
        elif bg_ok_count >= 4 and teal_found:
            print(f"PARTIAL: Component 7 — {bg_ok_count}/8 slides with #F5F5F5 bg, teal found (0.10 pts)")
            total_score += 0.10
        elif bg_ok_count >= 4 or teal_found:
            print(f"PARTIAL: Component 7 — bg count={bg_ok_count}, teal={teal_found} (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 7 — bg count={bg_ok_count}/8, teal accent={teal_found}")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entrypoint
if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(FILE_PATH)
