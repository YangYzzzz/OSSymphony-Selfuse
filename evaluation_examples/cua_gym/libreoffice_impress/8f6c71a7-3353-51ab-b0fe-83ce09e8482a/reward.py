"""
Reward Script: Create pie chart on slide 2 with budget allocation data and data labels
Task ID: impress_gf3_014
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Slide 2 contains a PIE chart
  Component 2 (0.30): Chart has correct 5 categories with correct percentage values
  Component 3 (0.25): Data labels show both category name and percentage
  Component 4 (0.20): Each slice has a distinct fill color
"""

import os
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_014'

# Expected categories and values
EXPECTED_CATEGORIES = ['Personnel', 'Equipment', 'Travel', 'Overhead', 'Miscellaneous']
EXPECTED_VALUES = [45.0, 20.0, 15.0, 12.0, 8.0]


def persist_app_state(domain):
    """Save any unsaved changes in LibreOffice."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for " + domain)
    except Exception as e:
        print("PERSIST_WARN: save hook failed: " + str(e))


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.chart import XL_CHART_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print("CRITICAL: Cannot load file " + file_path + ": " + str(e))
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 2 slides
    if len(prs.slides) < 2:
        print("FAIL: Presentation has fewer than 2 slides")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]

    # Find chart shape on slide 2
    chart_shape = None
    for shape in slide2.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Slide 2 contains a PIE chart (0.25 points)
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            if chart.chart_type == XL_CHART_TYPE.PIE:
                print("PASS: Component 1 — Slide 2 has a PIE chart (0.25 pts)")
                total_score += 0.25
            else:
                print("FAIL: Component 1 — Chart exists but type is " + str(chart.chart_type) + ", expected PIE")
        else:
            print("FAIL: Component 1 — No chart found on slide 2")
    except Exception as e:
        print("ERROR: Component 1 — " + str(e))

    # If no chart, remaining components cannot pass
    if chart_shape is None:
        print("\nScore: " + str(total_score) + "/1.0")
        print("REWARD: " + str(min(total_score, 1.0)))
        return min(total_score, 1.0)

    chart = chart_shape.chart

    # Component 2: Chart has correct 5 categories with correct values (0.30 points)
    try:
        plot = chart.plots[0]
        actual_categories = list(plot.categories)
        actual_values = list(chart.series[0].values)

        cats_match = len(actual_categories) == 5
        vals_match = len(actual_values) == 5

        if cats_match and vals_match:
            # Check each category-value pair
            matched_pairs = 0
            for exp_cat, exp_val in zip(EXPECTED_CATEGORIES, EXPECTED_VALUES):
                for act_cat, act_val in zip(actual_categories, actual_values):
                    if act_cat.strip().lower() == exp_cat.lower() and abs(float(act_val) - exp_val) < 0.5:
                        matched_pairs += 1
                        break

            if matched_pairs == 5:
                print("PASS: Component 2 — All 5 categories and values match (0.30 pts)")
                total_score += 0.30
            elif matched_pairs >= 3:
                partial = round(0.30 * matched_pairs / 5, 2)
                print("PARTIAL: Component 2 — " + str(matched_pairs) + "/5 pairs match (" + str(partial) + " pts)")
                total_score += partial
            else:
                print("FAIL: Component 2 — Only " + str(matched_pairs) + "/5 pairs match")
                print("  Expected categories: " + str(EXPECTED_CATEGORIES))
                print("  Actual categories: " + str(actual_categories))
                print("  Expected values: " + str(EXPECTED_VALUES))
                print("  Actual values: " + str(actual_values))
        else:
            print("FAIL: Component 2 — Expected 5 categories/values, got " + str(len(actual_categories)) + " categories, " + str(len(actual_values)) + " values")
    except Exception as e:
        print("ERROR: Component 2 — " + str(e))

    # Component 3: Data labels show both category name and percentage (0.25 points)
    try:
        plot = chart.plots[0]
        if plot.has_data_labels:
            dl = plot.data_labels
            has_cat_name = dl.show_category_name
            has_percentage = dl.show_percentage

            if has_cat_name and has_percentage:
                print("PASS: Component 3 — Data labels show category name AND percentage (0.25 pts)")
                total_score += 0.25
            elif has_cat_name or has_percentage:
                print("PARTIAL: Component 3 — Data labels show " +
                      ("category name" if has_cat_name else "") +
                      (" and " if has_cat_name and has_percentage else "") +
                      ("percentage" if has_percentage else "") +
                      " (0.125 pts)")
                total_score += 0.125
            else:
                print("FAIL: Component 3 — Data labels exist but show neither category name nor percentage")
        else:
            print("FAIL: Component 3 — No data labels on chart")
    except Exception as e:
        print("ERROR: Component 3 — " + str(e))

    # Component 4: Each slice has a distinct fill color (0.20 points)
    try:
        chart_xml = chart._chartSpace.xml
        ns = {
            'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        }
        root = ET.fromstring(chart_xml)
        points = root.findall('.//c:ser/c:dPt', ns)

        colors = set()
        for pt in points:
            fill = pt.find('.//a:solidFill/a:srgbClr', ns)
            if fill is not None:
                colors.add(fill.get('val'))

        if len(colors) >= 5:
            print("PASS: Component 4 — " + str(len(colors)) + " distinct fill colors found (0.20 pts)")
            total_score += 0.20
        elif len(colors) >= 3:
            partial = round(0.20 * len(colors) / 5, 2)
            print("PARTIAL: Component 4 — " + str(len(colors)) + " distinct colors, expected 5 (" + str(partial) + " pts)")
            total_score += partial
        else:
            print("FAIL: Component 4 — Only " + str(len(colors)) + " distinct fill colors found, expected 5")
    except Exception as e:
        print("ERROR: Component 4 — " + str(e))

    final_score = min(total_score, 1.0)
    print("\nScore: " + str(total_score) + "/1.0")
    print("REWARD: " + str(final_score))
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = WORKDIR + '/' + TASK_ID + '.pptx'
if not os.path.exists(file_path):
    print("File not found: " + file_path)
    print("REWARD: 0.0")
else:
    verify_task(file_path)
