"""
Initial Setup: Create a 5-slide Budget Overview presentation with empty Slide 2 content area
Task ID: impress_gf3_014
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_014'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Budget Overview"
    slide1.placeholders[1].text = "Fiscal Year 2025-2026\nPrepared by Finance Department"

    # --- Slide 2: Budget Allocation (title only, empty content area - NO chart) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Budget Allocation"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x62)

    # --- Slide 3: Revenue Streams ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox3_title = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf3 = txBox3_title.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Revenue Streams"
    p3.alignment = PP_ALIGN.LEFT
    r3 = p3.runs[0]
    r3.font.size = Pt(36)
    r3.font.bold = True
    r3.font.color.rgb = RGBColor(0x2E, 0x4A, 0x62)

    txBox3_body = slide3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(5))
    tf3b = txBox3_body.text_frame
    tf3b.word_wrap = True
    items = [
        ("Product Sales", "$2.4M", "Core product line contributing 48% of total revenue"),
        ("Service Contracts", "$1.1M", "Annual maintenance and support agreements"),
        ("Consulting", "$680K", "Professional services and implementation support"),
        ("Licensing", "$420K", "IP licensing and partnership royalties"),
    ]
    for i, (name, amount, desc) in enumerate(items):
        para = tf3b.paragraphs[0] if i == 0 else tf3b.add_paragraph()
        run_name = para.add_run()
        run_name.text = f"{name}: {amount}"
        run_name.font.size = Pt(20)
        run_name.font.bold = True
        run_name.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        para_desc = tf3b.add_paragraph()
        run_desc = para_desc.add_run()
        run_desc.text = f"   {desc}"
        run_desc.font.size = Pt(16)
        run_desc.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # --- Slide 4: Project Timeline ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox4_title = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf4 = txBox4_title.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Project Timeline"
    p4.alignment = PP_ALIGN.LEFT
    r4 = p4.runs[0]
    r4.font.size = Pt(36)
    r4.font.bold = True
    r4.font.color.rgb = RGBColor(0x2E, 0x4A, 0x62)

    # Add a simple table for timeline
    table_shape = slide4.shapes.add_table(5, 3, Inches(0.8), Inches(1.5), Inches(11), Inches(4))
    table = table_shape.table
    headers = ["Phase", "Duration", "Key Deliverables"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(16)

    timeline_data = [
        ["Planning & Analysis", "Q1 2025", "Requirements doc, stakeholder sign-off"],
        ["Development", "Q2-Q3 2025", "Core features, integration testing"],
        ["Pilot Rollout", "Q4 2025", "Beta deployment, user feedback collection"],
        ["Full Launch", "Q1 2026", "Company-wide deployment, training"],
    ]
    for r, row_data in enumerate(timeline_data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 5: Summary ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox5_title = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf5 = txBox5_title.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Summary & Next Steps"
    p5.alignment = PP_ALIGN.LEFT
    r5 = p5.runs[0]
    r5.font.size = Pt(36)
    r5.font.bold = True
    r5.font.color.rgb = RGBColor(0x2E, 0x4A, 0x62)

    txBox5_body = slide5.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(5))
    tf5b = txBox5_body.text_frame
    tf5b.word_wrap = True
    summary_items = [
        "Total budget allocation of $4.6M across five key categories",
        "Personnel costs represent the largest share at 45% of total budget",
        "Equipment modernization program requires 20% allocation",
        "Travel and conference budget maintained at 15% for team development",
        "Overhead and miscellaneous costs optimized to 20% combined",
        "Quarterly budget reviews scheduled starting Q2 2025",
    ]
    for i, item in enumerate(summary_items):
        para = tf5b.paragraphs[0] if i == 0 else tf5b.add_paragraph()
        run_item = para.add_run()
        run_item.text = f"• {item}"
        run_item.font.size = Pt(18)
        run_item.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
