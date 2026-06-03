"""
FINAL REWARD SCRIPT - SUCCESS
Task: I'm updating slide 9 with an integral formula for my math lecture—how do I include the equation ∫₀^∞ e^−x² dx = √π⁄2 using Math in LibreOffice Impress?
Generated: 2025-08-07 13:07:54
Status: success
Model: o4-mini
Total Steps: 3
"""

import os, re
from pptx import Presentation

def verify_task(file_path):
    print("Checking task completion...")
    total_score = 0.0
    max_score = 1.0

    # Requirement 1: File exists (0.2 points)
    try:
        if os.path.exists(file_path):
            print(f"✓ File exists: {file_path} (0.2 points)")
            total_score += 0.2
        else:
            print(f"✗ File not found: {file_path}")
            print(f"REWARD: {total_score}")
            return total_score
    except Exception as e:
        print(f"✗ Error checking file existence: {e}")
        print(f"REWARD: {total_score}")
        return total_score

    # Requirement 2: Load presentation and check slide count >=9 (0.3 points)
    try:
        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        print(f"✓ Presentation loaded with {slide_count} slides")
        if slide_count >= 9:
            print(f"✓ Slide count >= 9 (0.3 points)")
            total_score += 0.3
        else:
            print(f"✗ Slide count < 9: {slide_count} (0 points)")
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        print(f"REWARD: {total_score}")
        return total_score

    # Requirement 3: Check for integral formula text on slide 9 (0.5 points)
    try:
        slide9 = prs.slides[8]
        found_formula = False
        required_tokens = ['∫', 'dx', '=', '√', 'π']
        for shape in slide9.shapes:
            if hasattr(shape, 'text') and shape.text and shape.text.strip():
                original_text = shape.text.strip()
                # Normalize text: remove spaces and unify characters
                norm_text = original_text.replace(' ', '').replace('−', '-').replace('\u2082', '2')
                print(f"  Inspecting text: '{original_text}' => normalized: '{norm_text}'")
                if all(tok in norm_text for tok in required_tokens):
                    print(f"✓ Found integral formula in slide 9: '{original_text}' (0.5 points)")
                    total_score += 0.5
                    found_formula = True
                    break
        if not found_formula:
            print("✗ Integral formula not found on slide 9 (0 points)")
    except Exception as e:
        print(f"✗ Error checking slide 9 content: {e}")
        print(f"REWARD: {total_score}")
        return total_score

    # Final scoring
    final_score = min(total_score, max_score)
    print(f"Total score: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score

# Execute verification
if __name__ == '__main__':
    verify_task('/home/user/im_updating_slide_9_with_an_integral_formula_for_my_math_lecturehow_do_i_include_the_equation_₀_ex²_.pptx')
