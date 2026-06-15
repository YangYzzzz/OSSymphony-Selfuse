"""
Initial Setup: Create a Formatting Demo presentation with 8 slides.
Slide 6 has a title but no text boxes below it (empty content area).
Task ID: impress_gf3_041
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_041'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Formatting Demo"
    slide1.placeholders[1].text = "Advanced Text & Layout Techniques"

    # --- Slide 2: Introduction ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Introduction"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "This presentation demonstrates various formatting techniques available in modern presentation software."
    p2 = body2.add_paragraph()
    p2.text = "Topics include font styles, paragraph spacing, alignment, and text effects."
    p3 = body2.add_paragraph()
    p3.text = "Each slide focuses on a specific formatting category."

    # --- Slide 3: Font Styles ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Font Styles Overview"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Common font properties:"
    items3 = ["Bold and italic for emphasis", "Underline for hyperlinks and key terms",
              "Strikethrough for deleted content", "Font color for visual hierarchy"]
    for item in items3:
        p = body3.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 4: Alignment Techniques ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Text Alignment"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Left alignment is the default for most Western languages."
    p4a = body4.add_paragraph()
    p4a.text = "Center alignment works well for titles and headings."
    p4b = body4.add_paragraph()
    p4b.text = "Right alignment is useful for dates and numerical data."
    p4c = body4.add_paragraph()
    p4c.text = "Justified text creates clean edges on both sides of a text block."

    # --- Slide 5: Color and Themes ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Color & Theme Usage"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Consistent color usage reinforces brand identity."
    p5a = body5.add_paragraph()
    p5a.text = "Primary colors draw attention to key information."
    p5b = body5.add_paragraph()
    p5b.text = "Accent colors should be used sparingly for emphasis."
    p5c = body5.add_paragraph()
    p5c.text = "High contrast between text and background improves readability."

    # --- Slide 6: Paragraph Spacing (TASK TARGET) ---
    # Only a title, NO text boxes below - the agent must create the text box
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add just a title text box at the top
    title_box = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf6 = title_box.text_frame
    p6 = tf6.paragraphs[0]
    p6.text = "Paragraph Spacing"
    p6.alignment = PP_ALIGN.LEFT
    run6 = p6.runs[0]
    run6.font.size = Pt(28)
    run6.font.bold = True
    run6.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # --- Slide 7: Bullet Points ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Bullet Point Formatting"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "Project milestones for Q2 2025:"
    items7 = [
        ("Complete design review", 1),
        ("Finalize vendor contracts", 1),
        ("Budget approval from finance", 2),
        ("Stakeholder sign-off", 2),
        ("Launch beta testing phase", 1),
    ]
    for text, level in items7:
        p = body7.add_paragraph()
        p.text = text
        p.level = level

    # --- Slide 8: Summary ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Summary & Next Steps"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Key formatting principles covered in this session:"
    items8 = ["Font styles create visual hierarchy",
              "Alignment guides the reader's eye",
              "Spacing improves readability and flow",
              "Color reinforces meaning and brand"]
    for item in items8:
        p = body8.add_paragraph()
        p.text = item
        p.level = 1

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
