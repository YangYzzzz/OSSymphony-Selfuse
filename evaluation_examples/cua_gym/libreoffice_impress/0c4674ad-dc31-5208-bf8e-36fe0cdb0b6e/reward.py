"""
Reward Script: Create text box on slide 6 with line spacing and paragraph spacing
Task ID: impress_gf3_041
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Text box with 2 non-empty paragraphs exists on slide 6
  Component 2 (0.25): First paragraph has 1.5x line spacing (spcPct=150000)
  Component 3 (0.25): Second paragraph has 2.0x line spacing (spcPct=200000)
  Component 4 (0.25): Both paragraphs have space-after = 12pt (spcPts=1200)
"""

import os
from pptx import Presentation
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_041'


def get_line_spacing_pct(para):
    """Return line spacing as spcPct val (e.g. 150000 for 1.5x), or None if not set."""
    pPr = para._pPr
    if pPr is None:
        return None
    lnSpc = pPr.find(qn('a:lnSpc'))
    if lnSpc is None:
        return None
    spcPct = lnSpc.find(qn('a:spcPct'))
    if spcPct is not None:
        try:
            return int(spcPct.get('val'))
        except (TypeError, ValueError):
            return None
    return None


def get_space_after_pts(para):
    """Return space-after as spcPts val (e.g. 1200 for 12pt), or None if not set."""
    pPr = para._pPr
    if pPr is None:
        return None
    spcAft = pPr.find(qn('a:spcAft'))
    if spcAft is None:
        return None
    spcPts = spcAft.find(qn('a:spcPts'))
    if spcPts is not None:
        try:
            return int(spcPts.get('val'))
        except (TypeError, ValueError):
            return None
    return None


def find_task_textbox(slide):
    """Find a text box on slide 6 that has at least 2 non-empty paragraphs.
    Excludes the title placeholder and any pre-existing header textbox."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # Skip placeholders (title, etc.)
        if shape.shape_type == 14:  # PLACEHOLDER
            continue
        tf = shape.text_frame
        nonempty_paras = [p for p in tf.paragraphs if p.text.strip()]
        if len(nonempty_paras) >= 2:
            return shape
    return None


def verify_task(file_path):
    """Verify task completion with progressive scoring. Returns float 0.0-1.0."""
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[5]  # slide 6 (0-indexed)

    # Find the task-introduced text box (with 2+ non-empty paragraphs)
    task_shape = find_task_textbox(slide)

    # Component 1: Text box with 2 non-empty paragraphs exists on slide 6 (0.25 points)
    try:
        if task_shape is not None:
            nonempty_paras = [p for p in task_shape.text_frame.paragraphs if p.text.strip()]
            print(f"PASS: Component 1 -- Text box found with {len(nonempty_paras)} non-empty paragraphs (0.25 pts)")
            total_score += 0.25
        else:
            print("FAIL: Component 1 -- No text box with 2+ non-empty paragraphs found on slide 6")
            # No text box means all other components also fail
            print(f"\nScore: {total_score}/1.0")
            print(f"REWARD: {total_score}")
            return total_score
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Get the non-empty paragraphs for further checks
    nonempty_paras = [p for p in task_shape.text_frame.paragraphs if p.text.strip()]

    # Component 2: First paragraph has 1.5x line spacing (0.25 points)
    try:
        para1 = nonempty_paras[0]
        ls1 = get_line_spacing_pct(para1)
        if ls1 is not None and ls1 == 150000:
            print(f"PASS: Component 2 -- First paragraph line spacing = 150000 (1.5x) (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 -- First paragraph line spacing = {ls1}, expected 150000 (1.5x)")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Second paragraph has 2.0x line spacing (0.25 points)
    try:
        para2 = nonempty_paras[1]
        ls2 = get_line_spacing_pct(para2)
        if ls2 is not None and ls2 == 200000:
            print(f"PASS: Component 3 -- Second paragraph line spacing = 200000 (2.0x) (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 -- Second paragraph line spacing = {ls2}, expected 200000 (2.0x)")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Both paragraphs have space-after = 12pt (0.25 points)
    try:
        sa1 = get_space_after_pts(nonempty_paras[0])
        sa2 = get_space_after_pts(nonempty_paras[1])
        if sa1 == 1200 and sa2 == 1200:
            print(f"PASS: Component 4 -- Both paragraphs have space-after = 1200 (12pt) (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 4 -- Space-after: para1={sa1}, para2={sa2}, expected 1200 for both")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook: save any unsaved LibreOffice state
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
