"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 204’s heading is stuck to the left and still has an underline. In LibreOffice Impress, how can I bump that title to the exact center and strip the underline completely?
Generated: 2025-09-10 19:34:44
Status: success
Model: azure-o3
Total Steps: 2
"""

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
import os

def verify_slide_204_title(file_path: str) -> float:
    """Verify that on slide 204 the title is centered and not underlined.

    Scoring (progressive):
        0.5 points – Every non-empty title paragraph is CENTER aligned
        0.5 points – No text run in the title has any underline applied
    Returns a float between 0.0 and 1.0 (inclusive).
    """
    print(f"Verifying slide 204 title in: {file_path}\n")

    total_score = 0.0
    max_score   = 1.0

    # ---------- Basic file & slide existence checks (no points) ----------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        return 0.0

    slide_index = 203  # zero-based index for slide 204
    if len(prs.slides) <= slide_index:
        print(f"✗ Slide 204 is missing (presentation only has {len(prs.slides)} slides)")
        return 0.0

    slide = prs.slides[slide_index]
    title_shape = slide.shapes.title
    if title_shape is None:
        print("✗ No title shape found on slide 204")
        return 0.0

    # ---------- Requirement 1: Title paragraphs are centered ----------
    tf = title_shape.text_frame
    paragraphs = [p for p in tf.paragraphs if (p.text or '').strip()]

    if not paragraphs:
        print("✗ Title contains no text; cannot verify alignment")
    else:
        alignment_ok = all(p.alignment == PP_ALIGN.CENTER for p in paragraphs)
        for p in paragraphs:
            print(f"Paragraph '{p.text}' alignment value: {p.alignment}")
        if alignment_ok:
            print("✓ All non-empty paragraphs are centered (0.5 points)")
            total_score += 0.5
        else:
            print("✗ Not all title paragraphs are centered (0 points)")

    # ---------- Requirement 2: No underline on any run ----------
    underline_present = False
    for p in paragraphs:
        for run in p.runs:
            val = run.font.underline
            print(f"Run text '{run.text}' underline value: {val}")
            # pptx returns True / False / None / MSO_THEME_COLOR / enum values etc.
            if val is True:
                underline_present = True
            elif val not in [None, False]:
                # Any explicit underline style other than None/False is unwanted
                underline_present = True
    if not underline_present:
        print("✓ No underline detected on any run (0.5 points)")
        total_score += 0.5
    else:
        print("✗ Underline still present on one or more runs (0 points)")

    # ---------- Final score ----------
    final_score = min(total_score, max_score)
    print(f"\nTotal score: {final_score}\n")
    return final_score


if __name__ == "__main__":
    FILE_PATH = "/home/user/slide_204s_heading_is_stuck_to_the_left_and_still_has_an_underline_in_libreoffice_impress_how_can_i__golden.pptx"
    reward = verify_slide_204_title(FILE_PATH)
    print(f"REWARD: {reward}")
