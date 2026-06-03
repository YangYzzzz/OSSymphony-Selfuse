"""
Initial Setup: Create a 5-slide presentation with 'Quarterly Overview' title on slide 1
Task ID: impress_tct_002
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_002'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title + Content layout - "Quarterly Overview" ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide1.shapes.title.text = "Quarterly Overview"
    # Leave content area empty (placeholder 1 exists but no text added)

    # --- Slide 2: Regional Performance ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Regional Performance"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "North America: Strong growth in enterprise segment"
    p2 = body2.add_paragraph()
    p2.text = "Europe: Steady expansion in mid-market accounts"
    p2.level = 0
    p3 = body2.add_paragraph()
    p3.text = "Asia-Pacific: Emerging opportunities in healthcare vertical"
    p3.level = 0

    # --- Slide 3: Key Metrics ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Key Metrics"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Customer Retention Rate: 94.2%"
    for metric in [
        "Net Promoter Score: 72",
        "Average Deal Size: $48,500",
        "Pipeline Coverage Ratio: 3.1x",
        "Win Rate: 38%",
    ]:
        p = body3.add_paragraph()
        p.text = metric
        p.level = 0

    # --- Slide 4: Strategic Initiatives ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Strategic Initiatives"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Launch AI-powered analytics dashboard by Q3"
    for item in [
        "Expand partner ecosystem in Latin America",
        "Implement automated onboarding workflow",
        "Migrate legacy clients to cloud platform",
    ]:
        p = body4.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 5: Next Steps ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Next Steps"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Review updated pricing model with finance team"
    for item in [
        "Schedule quarterly business reviews with top 20 accounts",
        "Finalize hiring plan for Q3 engineering sprint",
    ]:
        p = body5.add_paragraph()
        p.text = item
        p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
