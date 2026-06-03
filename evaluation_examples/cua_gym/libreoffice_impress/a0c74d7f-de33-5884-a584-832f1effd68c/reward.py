"""
Reward Script: Create a 3x3 table on slide 1 with specific headers
Task ID: impress_tct_002
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): Table exists on slide 1
  Component 2 (0.3): Table is 3 rows x 3 columns
  Component 3 (0.4): Header row contains 'Product', 'Q1 Sales', 'Q2 Sales'
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_002'


def persist_app_state(domain: str):
    """Save any unsaved GUI edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 1 slide
    if len(prs.slides) < 1:
        print("FAIL: Presentation has no slides")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]

    # Find table(s) on slide 1
    tables = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            tables.append(shape.table)

    # Component 1: A table exists on slide 1 (0.3 points)
    try:
        if len(tables) > 0:
            print(f"PASS: Component 1 - Table found on slide 1 ({len(tables)} table(s)) (0.3 pts)")
            total_score += 0.3
        else:
            print("FAIL: Component 1 - No table found on slide 1")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    if len(tables) == 0:
        # No table means no further checks possible
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Use the first table found
    table = tables[0]

    # Component 2: Table dimensions are 3 rows x 3 columns (0.3 points)
    try:
        num_rows = len(table.rows)
        num_cols = len(table.columns)
        if num_rows == 3 and num_cols == 3:
            print(f"PASS: Component 2 - Table is 3x3 (rows={num_rows}, cols={num_cols}) (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 - Expected 3x3 table, found {num_rows}x{num_cols}")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Header row contains 'Product', 'Q1 Sales', 'Q2 Sales' (0.4 points)
    try:
        expected_headers = ['Product', 'Q1 Sales', 'Q2 Sales']
        # Check each header cell; award partial credit per correct header
        correct_count = 0
        for col_idx, expected in enumerate(expected_headers):
            if col_idx < len(table.columns):
                actual = table.cell(0, col_idx).text.strip()
                if actual == expected:
                    correct_count += 1
                    print(f"  Header [{col_idx}]: '{actual}' == '{expected}' OK")
                else:
                    print(f"  Header [{col_idx}]: '{actual}' != '{expected}' MISMATCH")
            else:
                print(f"  Header [{col_idx}]: column does not exist")

        if correct_count == 3:
            print(f"PASS: Component 3 - All 3 headers correct (0.4 pts)")
            total_score += 0.4
        elif correct_count > 0:
            partial = round(0.4 * correct_count / 3, 2)
            print(f"PARTIAL: Component 3 - {correct_count}/3 headers correct ({partial} pts)")
            total_score += partial
        else:
            print("FAIL: Component 3 - No headers match")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
