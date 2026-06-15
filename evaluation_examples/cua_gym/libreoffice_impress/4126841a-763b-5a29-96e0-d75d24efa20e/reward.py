"""
Reward Script: Waterfall chart on slide 5 showing profit bridge
Task ID: impress_exec_036
Domain: libreoffice_impress
Scoring:
  Component 1: Chart exists on slide 5 (0.25)
  Component 2: Chart title is 'Profit Bridge Analysis' (0.15)
  Component 3: Categories match 7 profit bridge items (0.25)
  Component 4: Chart data values match expected waterfall data (0.25)
  Component 5: Slide title text updated to 'Profit Bridge Analysis' (0.10)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_036'


def persist_app_state(domain: str):
    """Best-effort save of any open LibreOffice document."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
    except ImportError:
        print("CRITICAL: python-pptx not installed")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 5 slides
    if len(prs.slides) < 5:
        print(f"PRECONDITION FAIL: Need >= 5 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # slide 5 (0-indexed)

    # Find chart shape on slide 5
    chart_shape = None
    for shape in slide.shapes:
        if shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Chart exists on slide 5 (0.25 points)
    try:
        if chart_shape is not None:
            print(f"PASS: Component 1 - Chart found on slide 5 (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 - No chart found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # If no chart, remaining components cannot pass
    if chart_shape is None:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    chart = chart_shape.chart

    # Component 2: Chart title is 'Profit Bridge Analysis' (0.15 points)
    try:
        if chart.has_title:
            title_text = ""
            if chart.chart_title.has_text_frame:
                title_text = chart.chart_title.text_frame.text.strip()
            if title_text == "Profit Bridge Analysis":
                print(f"PASS: Component 2 - Chart title is 'Profit Bridge Analysis' (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 2 - Chart title is '{title_text}', expected 'Profit Bridge Analysis'")
        else:
            print(f"FAIL: Component 2 - Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Categories match the 7 profit bridge items (0.25 points)
    # Expected categories (may contain line breaks in golden but core labels must match)
    expected_labels = [
        "Starting Revenue", "COGS", "Gross Profit", "OpEx", "EBITDA", "D&A", "Net Income"
    ]
    try:
        raw_cats = [str(c) for c in chart.plots[0].categories]
        # Normalize: replace newlines with spaces and strip
        actual_cats = [c.replace('\n', ' ').strip() for c in raw_cats]

        matched = 0
        for exp in expected_labels:
            # Check if any actual category matches (case-insensitive, whitespace-normalized)
            exp_norm = exp.lower().replace(' ', '')
            for act in actual_cats:
                act_norm = act.lower().replace(' ', '')
                if exp_norm == act_norm:
                    matched += 1
                    break

        if matched == len(expected_labels):
            print(f"PASS: Component 3 - All 7 categories present: {actual_cats} (0.25 pts)")
            total_score += 0.25
        elif matched >= 5:
            partial = 0.25 * (matched / len(expected_labels))
            print(f"PARTIAL: Component 3 - {matched}/7 categories matched ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 - Only {matched}/7 categories matched. Found: {actual_cats}")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Chart data values match expected waterfall structure (0.25 points)
    # The task specifies: Starting Revenue ($50M), COGS (-$20M), Gross Profit ($30M),
    # OpEx (-$18M), EBITDA ($12M), D&A (-$3M), Net Income ($9M)
    # In a waterfall stacked bar, the visible values should be: 50, 20, 30, 18, 12, 3, 9
    expected_visible_values = {
        "Starting Revenue": 50.0,
        "COGS": 20.0,
        "Gross Profit": 30.0,
        "OpEx": 18.0,
        "EBITDA": 12.0,
        "D&A": 3.0,
        "Net Income": 9.0,
    }
    try:
        num_series = len(chart.series)
        # Collect all values across series per category
        all_values = []
        for s_idx in range(num_series):
            vals = list(chart.series[s_idx].values)
            all_values.append(vals)

        num_cats = len(list(chart.plots[0].categories))

        # For each category, sum the non-zero visible values
        # In a waterfall chart, the visible bar height should match expected values
        value_matches = 0
        total_checks = 7

        for cat_idx in range(min(num_cats, 7)):
            # Collect all non-zero values for this category across all series
            cat_values = []
            for s_idx in range(num_series):
                if s_idx < len(all_values) and cat_idx < len(all_values[s_idx]):
                    v = all_values[s_idx][cat_idx]
                    if v is not None and abs(v) > 0.01:
                        cat_values.append(v)

            exp_label = expected_labels[cat_idx] if cat_idx < len(expected_labels) else f"cat_{cat_idx}"
            exp_val = expected_visible_values.get(exp_label, None)

            if exp_val is not None:
                # Check if any combination of visible values matches
                # For totals/positive items: single visible value == expected
                # For negative items: visible value == expected (absolute)
                found_match = False
                for v in cat_values:
                    if abs(v - exp_val) < 0.5:
                        found_match = True
                        break

                if found_match:
                    value_matches += 1
                else:
                    print(f"  Data mismatch at '{exp_label}': expected ~{exp_val}, found values {cat_values}")

        if value_matches == total_checks:
            print(f"PASS: Component 4 - All 7 data values match expected waterfall amounts (0.25 pts)")
            total_score += 0.25
        elif value_matches >= 4:
            partial = 0.25 * (value_matches / total_checks)
            print(f"PARTIAL: Component 4 - {value_matches}/7 data values match ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 - Only {value_matches}/7 data values match")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Slide title/heading text updated to 'Profit Bridge Analysis' (0.10 points)
    # In the initial file, the text box says "Profit Bridge". In golden, it says "Profit Bridge Analysis".
    try:
        title_found = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = para.text.strip()
                    if txt == "Profit Bridge Analysis":
                        title_found = True
                        break
            if title_found:
                break

        if title_found:
            print(f"PASS: Component 5 - Slide text 'Profit Bridge Analysis' found (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 5 - Slide text 'Profit Bridge Analysis' not found")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
