"""
Initial Setup: Create a 10-slide CFO Deep Dive presentation with slide 5 titled 'Profit Bridge' (no chart).
Task ID: impress_exec_036
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_036'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                color=None, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_slide_with_title(prs, layout_idx, title_text, subtitle_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    if subtitle_text and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle_text
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "CFO Deep Dive"
    slide1.placeholders[1].text = "Quarterly Financial Review - Q1 2025\nPrepared by Finance Team"

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_textbox(slide2, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                "Executive Summary", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x6B))
    summary_items = [
        "Revenue grew 12% YoY to $50M, driven by enterprise segment expansion",
        "Gross margin improved 200bps to 60% through supply chain optimization",
        "Operating expenses held flat at $18M despite 15% headcount growth",
        "Free cash flow generation of $7.2M, up from $5.1M in Q4 2024",
        "Net income reached $9M, a 22% improvement over prior quarter"
    ]
    for i, item in enumerate(summary_items):
        add_textbox(slide2, Inches(0.8), Inches(1.5 + i * 0.9), Inches(11), Inches(0.8),
                    f"  {item}", font_size=16)

    # --- Slide 3: Revenue Breakdown ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide3, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                "Revenue Breakdown by Segment", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x6B))
    # Table with revenue data
    table_shape = slide3.shapes.add_table(6, 4, Inches(1), Inches(1.5), Inches(10), Inches(3.5))
    table = table_shape.table
    headers = ["Segment", "Q1 2025 ($M)", "Q4 2024 ($M)", "YoY Growth"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    data = [
        ["Enterprise", "$22.5", "$19.8", "+13.6%"],
        ["Mid-Market", "$14.0", "$13.2", "+6.1%"],
        ["SMB", "$8.5", "$7.5", "+13.3%"],
        ["Professional Services", "$3.5", "$3.0", "+16.7%"],
        ["Other", "$1.5", "$1.1", "+36.4%"],
    ]
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 4: Cost Structure ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide4, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                "Cost Structure Analysis", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x6B))
    cost_items = [
        ("Cost of Goods Sold (COGS)", "$20.0M", "40% of revenue"),
        ("Research & Development", "$7.2M", "14.4% of revenue"),
        ("Sales & Marketing", "$6.5M", "13.0% of revenue"),
        ("General & Administrative", "$4.3M", "8.6% of revenue"),
        ("Depreciation & Amortization", "$3.0M", "6.0% of revenue"),
    ]
    for i, (name, amount, pct) in enumerate(cost_items):
        y = Inches(1.5 + i * 1.0)
        add_textbox(slide4, Inches(1), y, Inches(5), Inches(0.5),
                    name, font_size=16, bold=True)
        add_textbox(slide4, Inches(7), y, Inches(2), Inches(0.5),
                    amount, font_size=16, bold=True,
                    color=RGBColor(0xC0, 0x39, 0x2B))
        add_textbox(slide4, Inches(9.5), y, Inches(3), Inches(0.5),
                    pct, font_size=14)

    # --- Slide 5: Profit Bridge (NO CHART - this is what the agent must add) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide5, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                "Profit Bridge", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x6B))
    add_textbox(slide5, Inches(2), Inches(3), Inches(8), Inches(1.5),
                "[ Chart placeholder - Waterfall chart to be added ]",
                font_size=18, color=RGBColor(0x99, 0x99, 0x99),
                alignment=PP_ALIGN.CENTER)

    # --- Slide 6: Cash Flow Statement ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide6, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                "Cash Flow Statement", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x6B))
    cf_table = slide6.shapes.add_table(5, 3, Inches(1.5), Inches(1.5), Inches(9), Inches(3))
    cf = cf_table.table
    cf_headers = ["Category", "Q1 2025 ($M)", "Q4 2024 ($M)"]
    for i, h in enumerate(cf_headers):
        cell = cf.cell(0, i)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
    cf_data = [
        ["Operating Cash Flow", "$10.8", "$9.2"],
        ["Investing Cash Flow", "($2.1)", "($3.5)"],
        ["Financing Cash Flow", "($1.5)", "($0.6)"],
        ["Net Change in Cash", "$7.2", "$5.1"],
    ]
    for r, row_data in enumerate(cf_data, 1):
        for c, val in enumerate(row_data):
            cf.cell(r, c).text = val

    # --- Slide 7: Balance Sheet Highlights ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide7, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                "Balance Sheet Highlights", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x6B))
    bs_items = [
        "Total Assets: $185.4M (up from $172.1M)",
        "Cash & Equivalents: $42.3M",
        "Accounts Receivable: $28.7M (DSO: 52 days)",
        "Total Debt: $35.0M (Debt/EBITDA: 2.9x)",
        "Shareholders' Equity: $98.2M",
    ]
    for i, item in enumerate(bs_items):
        add_textbox(slide7, Inches(1), Inches(1.5 + i * 0.9), Inches(10), Inches(0.7),
                    f"  {item}", font_size=16)

    # --- Slide 8: Key Metrics Dashboard ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide8, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                "Key Performance Metrics", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x6B))
    metrics = [
        ("Revenue Growth", "12.0%", "Target: 10%"),
        ("Gross Margin", "60.0%", "Target: 58%"),
        ("EBITDA Margin", "24.0%", "Target: 22%"),
        ("Customer Retention", "94.2%", "Target: 93%"),
        ("ARR per Employee", "$285K", "Target: $270K"),
    ]
    for i, (name, val, target) in enumerate(metrics):
        y = Inches(1.5 + i * 1.0)
        add_textbox(slide8, Inches(1), y, Inches(4), Inches(0.5),
                    name, font_size=18, bold=True)
        add_textbox(slide8, Inches(5.5), y, Inches(2.5), Inches(0.5),
                    val, font_size=20, bold=True,
                    color=RGBColor(0x27, 0xAE, 0x60))
        add_textbox(slide8, Inches(8.5), y, Inches(3), Inches(0.5),
                    target, font_size=14, color=RGBColor(0x77, 0x77, 0x77))

    # --- Slide 9: Risk Factors ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide9, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                "Risk Factors & Mitigation", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x6B))
    risks = [
        "FX Exposure: 18% of revenue in EUR/GBP - hedging program covers 75%",
        "Customer Concentration: Top 5 clients = 32% of revenue (down from 38%)",
        "Talent Retention: Engineering attrition at 8.5% vs industry 12%",
        "Supply Chain: Dual-sourced 85% of critical components",
    ]
    for i, risk in enumerate(risks):
        add_textbox(slide9, Inches(0.8), Inches(1.5 + i * 1.2), Inches(11), Inches(0.9),
                    f"  {risk}", font_size=15)

    # --- Slide 10: Outlook & Guidance ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide10, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                "FY2025 Outlook & Guidance", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x6B))
    outlook_items = [
        "Revenue Guidance: $205M - $215M (10-15% growth)",
        "EBITDA Margin Target: 23-25%",
        "CapEx Budget: $12M (focused on platform infrastructure)",
        "Headcount Plan: +45 FTEs (primarily Engineering & Sales)",
        "M&A Pipeline: 2 strategic targets under evaluation",
    ]
    for i, item in enumerate(outlook_items):
        add_textbox(slide10, Inches(1), Inches(1.5 + i * 0.9), Inches(10), Inches(0.7),
                    f"  {item}", font_size=16)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
