"""
Reward Script: Bold all text on all slides and underline only the title on slide 1.
Task ID: osworld_impress_bold_all_title_size_underline_009
Domain: libreoffice_impress
Scoring:
  Component 1 (0.6 pts): All text runs across all 5 slides have bold=True
  Component 2 (0.2 pts): Slide 1 title run is underlined (underline=True)
  Component 3 (0.2 pts): No text runs outside slide 1 title have underline=True
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_bold_all_title_size_underline_009'


def get_all_text_runs(slide):
    """
    Recursively collect all (shape_name, para_idx, run) tuples from a slide,
    including runs inside GROUP shapes.
    Only returns runs with non-empty text.
    """
    def extract(shape):
        results = []
        if hasattr(shape, 'text_frame') and shape.has_text_frame:
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                for run in para.runs:
                    if (run.text or '').strip():
                        results.append((shape.name, para_idx, run))
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.extend(extract(sub))
        return results

    all_runs = []
    for shape in slide.shapes:
        all_runs.extend(extract(shape))
    return all_runs


def get_title_shape(slide):
    """Return the Title placeholder shape of a slide, or None if not found."""
    for shape in slide.shapes:
        if shape.has_text_frame and 'Title' in shape.name and 'Content' not in shape.name:
            return shape
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    try:
        from pptx import Presentation
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Sanity check: must have 5 slides
    num_slides = len(prs.slides)
    if num_slides != 5:
        print(f"FAIL: Expected 5 slides, found {num_slides}. Cannot score.")
        print("REWARD: 0.0")
        return 0.0

    total_score = 0.0

    # ----------------------------------------------------------------
    # Component 1: All text runs across all 5 slides have bold=True
    # (0.6 points)
    # In the initial state, all runs have bold=False.
    # In the golden state, all runs have bold=True.
    # ----------------------------------------------------------------
    try:
        non_bold_details = []

        for slide_idx, slide in enumerate(prs.slides):
            runs = get_all_text_runs(slide)
            for shape_name, para_idx, run in runs:
                bold_val = run.font.bold
                # bold=True means explicitly bold; bold=None means inherit (treat as NOT explicitly bold)
                if bold_val is not True:
                    non_bold_details.append(
                        f"slide {slide_idx+1}, shape '{shape_name}', para[{para_idx}], "
                        f"text={repr(run.text[:30])}, bold={bold_val}"
                    )

        if len(non_bold_details) == 0:
            print("PASS: Component 1 — All text runs across all 5 slides are bold (0.6 pts)")
            total_score += 0.6
        else:
            print(f"FAIL: Component 1 — {len(non_bold_details)} run(s) are not bold:")
            for detail in non_bold_details[:5]:  # Show at most 5 examples
                print(f"  - {detail}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # ----------------------------------------------------------------
    # Component 2: Slide 1 title run is underlined (underline=True)
    # (0.2 points)
    # In the initial state, slide 1 title has underline=False.
    # In the golden state, slide 1 title has underline=True.
    # ----------------------------------------------------------------
    try:
        slide1 = prs.slides[0]
        title_shape = get_title_shape(slide1)

        if title_shape is None:
            print("FAIL: Component 2 — No title shape found on slide 1")
        else:
            title_runs = get_all_text_runs(slide1)
            # Filter to only runs from the title shape
            title_runs_in_title = [
                (sn, pi, r) for (sn, pi, r) in title_runs
                if sn == title_shape.name
            ]

            if not title_runs_in_title:
                print("FAIL: Component 2 — No runs found in slide 1 title shape")
            else:
                # All runs in the title placeholder must be underlined
                all_title_underlined = all(
                    r.font.underline is True
                    for (_, _, r) in title_runs_in_title
                )
                if all_title_underlined:
                    title_text = title_shape.text_frame.paragraphs[0].text if title_shape.text_frame.paragraphs else ""
                    print(f"PASS: Component 2 — Slide 1 title is underlined: {repr(title_text[:40])} (0.2 pts)")
                    total_score += 0.2
                else:
                    for (sn, pi, r) in title_runs_in_title:
                        print(
                            f"FAIL: Component 2 — Slide 1 title run '{repr(r.text[:30])}' "
                            f"has underline={r.font.underline}, expected True"
                        )
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # ----------------------------------------------------------------
    # Component 3: No text runs OUTSIDE slide 1 title have underline=True
    # (0.2 points)
    # In the initial state, no text is underlined → this component passes.
    # BUT since this is a precondition check for initial, we tie it to
    # Component 2's failure on initial: if bold (Component 1) fails,
    # we rely on Component 1 alone to drive the initial → 0.0 path.
    #
    # To ensure Component 3 alone does NOT contribute to initial scoring,
    # we only award these 0.2 points if Component 1 also passed
    # (i.e., all text is bold). This ties this component to task-introduced changes.
    # ----------------------------------------------------------------
    try:
        # Only evaluate if Component 1 passed (all text is bold — the primary change)
        if total_score >= 0.6:
            spurious_underlines = []

            for slide_idx, slide in enumerate(prs.slides):
                runs = get_all_text_runs(slide)
                title_shape_for_slide = get_title_shape(slide) if slide_idx == 0 else None

                for shape_name, para_idx, run in runs:
                    # Skip slide 1 title runs — those are expected to be underlined
                    if slide_idx == 0 and title_shape_for_slide and shape_name == title_shape_for_slide.name:
                        continue
                    # All other runs must NOT be underlined
                    if run.font.underline is True:
                        spurious_underlines.append(
                            f"slide {slide_idx+1}, shape '{shape_name}', para[{para_idx}], "
                            f"text={repr(run.text[:30])}"
                        )

            if not spurious_underlines:
                print("PASS: Component 3 — No spurious underlines outside slide 1 title (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 3 — {len(spurious_underlines)} run(s) are incorrectly underlined:")
                for detail in spurious_underlines[:5]:
                    print(f"  - {detail}")
        else:
            print("SKIP: Component 3 — Skipped because Component 1 (all bold) did not pass")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the VM environment
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
