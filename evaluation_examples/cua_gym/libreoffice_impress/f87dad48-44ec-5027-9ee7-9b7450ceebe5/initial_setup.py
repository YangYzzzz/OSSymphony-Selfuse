"""
Initial Setup: Bold all text on all slides and underline only the title on slide 1.
Task ID: osworld_impress_bold_all_title_size_underline_009
Domain: libreoffice_impress

Creates a 5-slide marketing pitch deck where all text is regular weight (not bold,
not underlined) — the pre-task state.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_bold_all_title_size_underline_009'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_run(para, text, font_size=None, bold=False, italic=False,
                 underline=False, color=None):
    """Add a run to a paragraph with the specified formatting."""
    run = para.add_run()
    run.text = text
    if font_size:
        run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.underline = underline
    if color:
        run.font.color.rgb = RGBColor(*color)
    return run


def set_placeholder_text(placeholder, text, font_size=None, bold=False,
                          italic=False, underline=False, color=None,
                          alignment=None):
    """Set placeholder text with specific formatting."""
    tf = placeholder.text_frame
    tf.clear()
    para = tf.paragraphs[0]
    if alignment:
        para.alignment = alignment
    add_text_run(para, text, font_size=font_size, bold=bold, italic=italic,
                 underline=underline, color=color)


def create_initial():
    prs = Presentation()

    # Standard widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title Slide ──────────────────────────────────────────────────
    layout_title = prs.slide_layouts[0]  # Title Slide layout
    slide1 = prs.slides.add_slide(layout_title)

    title1 = slide1.shapes.title
    title1.text = ""
    tf = title1.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "NovaSpark Solutions"
    run.font.size = Pt(40)
    run.font.bold = False
    run.font.underline = False
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)  # dark blue

    subtitle1 = slide1.placeholders[1]
    tf_s = subtitle1.text_frame
    tf_s.clear()
    p_s = tf_s.paragraphs[0]
    p_s.alignment = PP_ALIGN.CENTER
    r_s = p_s.add_run()
    r_s.text = "Powering the Next Generation of Enterprise AI"
    r_s.font.size = Pt(22)
    r_s.font.bold = False
    r_s.font.italic = False
    r_s.font.underline = False
    r_s.font.color.rgb = RGBColor(0x40, 0x40, 0x40)

    # ── Slide 2: Problem Statement ────────────────────────────────────────────
    layout_content = prs.slide_layouts[1]  # Title and Content
    slide2 = prs.slides.add_slide(layout_content)

    title2 = slide2.shapes.title
    tf2 = title2.text_frame
    tf2.clear()
    p2t = tf2.paragraphs[0]
    r2t = p2t.add_run()
    r2t.text = "The Challenge"
    r2t.font.size = Pt(32)
    r2t.font.bold = False
    r2t.font.underline = False
    r2t.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    content2 = slide2.placeholders[1]
    tf2c = content2.text_frame
    tf2c.clear()
    bullets2 = [
        "Enterprises face 3x more data than they can process manually",
        "Legacy workflows cost an average of $2.4M per year in inefficiencies",
        "Talent shortages leave 62% of AI initiatives stalled at proof-of-concept",
        "Siloed tools create friction and slow down time-to-insight by months",
    ]
    for i, bullet in enumerate(bullets2):
        if i == 0:
            para = tf2c.paragraphs[0]
        else:
            para = tf2c.add_paragraph()
        para.level = 0
        run = para.add_run()
        run.text = bullet
        run.font.size = Pt(18)
        run.font.bold = False
        run.font.underline = False

    # ── Slide 3: Our Solution ─────────────────────────────────────────────────
    slide3 = prs.slides.add_slide(layout_content)

    title3 = slide3.shapes.title
    tf3 = title3.text_frame
    tf3.clear()
    p3t = tf3.paragraphs[0]
    r3t = p3t.add_run()
    r3t.text = "Our Solution"
    r3t.font.size = Pt(32)
    r3t.font.bold = False
    r3t.font.underline = False
    r3t.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    content3 = slide3.placeholders[1]
    tf3c = content3.text_frame
    tf3c.clear()
    bullets3 = [
        "NovaSpark AI Platform: end-to-end automation for enterprise workflows",
        "Integrates seamlessly with existing ERP, CRM, and data warehouse systems",
        "Reduces time-to-insight from weeks to hours using real-time ML pipelines",
        "Scalable SaaS pricing — start at $4,999/month, no upfront hardware",
    ]
    for i, bullet in enumerate(bullets3):
        if i == 0:
            para = tf3c.paragraphs[0]
        else:
            para = tf3c.add_paragraph()
        para.level = 0
        run = para.add_run()
        run.text = bullet
        run.font.size = Pt(18)
        run.font.bold = False
        run.font.underline = False

    # ── Slide 4: Traction & Metrics ───────────────────────────────────────────
    slide4 = prs.slides.add_slide(layout_content)

    title4 = slide4.shapes.title
    tf4 = title4.text_frame
    tf4.clear()
    p4t = tf4.paragraphs[0]
    r4t = p4t.add_run()
    r4t.text = "Traction & Key Metrics"
    r4t.font.size = Pt(32)
    r4t.font.bold = False
    r4t.font.underline = False
    r4t.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    content4 = slide4.placeholders[1]
    tf4c = content4.text_frame
    tf4c.clear()
    metrics = [
        "127 enterprise clients across 14 industries in under 18 months",
        "ARR grew 340% year-over-year — from $1.2M to $5.3M",
        "Net Promoter Score (NPS) of 72 — well above industry average of 31",
        "3 Fortune 500 pilot agreements signed in Q1 2025",
    ]
    for i, metric in enumerate(metrics):
        if i == 0:
            para = tf4c.paragraphs[0]
        else:
            para = tf4c.add_paragraph()
        para.level = 0
        run = para.add_run()
        run.text = metric
        run.font.size = Pt(18)
        run.font.bold = False
        run.font.underline = False

    # ── Slide 5: Call to Action ───────────────────────────────────────────────
    layout_title_only = prs.slide_layouts[5]  # Blank layout
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])

    title5 = slide5.shapes.title
    tf5 = title5.text_frame
    tf5.clear()
    p5t = tf5.paragraphs[0]
    p5t.alignment = PP_ALIGN.CENTER
    r5t = p5t.add_run()
    r5t.text = "Join Us — Let's Build the Future"
    r5t.font.size = Pt(32)
    r5t.font.bold = False
    r5t.font.underline = False
    r5t.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    content5 = slide5.placeholders[1]
    tf5c = content5.text_frame
    tf5c.clear()
    cta_lines = [
        "Seeking $8M Series A to accelerate product development and global expansion",
        "Target markets: North America, Western Europe, Southeast Asia",
        "Contact: partnerships@novaspark.ai  |  +1 (415) 882-0340",
        "Schedule a demo: novaspark.ai/demo",
    ]
    for i, line in enumerate(cta_lines):
        if i == 0:
            para = tf5c.paragraphs[0]
        else:
            para = tf5c.add_paragraph()
        para.level = 0
        run = para.add_run()
        run.text = line
        run.font.size = Pt(18)
        run.font.bold = False
        run.font.underline = False

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
