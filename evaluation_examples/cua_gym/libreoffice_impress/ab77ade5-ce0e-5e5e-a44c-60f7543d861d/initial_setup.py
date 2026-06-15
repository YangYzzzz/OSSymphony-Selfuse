"""
initial_setup.py - Creates Legacy_Data.pptx with 5 slides.
Slide 4 has a 4-column by 9-row table. Last 3 rows contain deprecated entries.
"""
import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

OUTPUT_PATH = "/home/user/Legacy_Data.pptx"

prs = Presentation()

# ── Slide 1: Title Slide ──
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
slide1.shapes.title.text = "Quarterly Infrastructure Report"
slide1.placeholders[1].text = "Legacy Systems Division — Q3 2024 Review"

# ── Slide 2: Overview (Title + Content) ──
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
slide2.shapes.title.text = "Project Overview"
body2 = slide2.placeholders[1].text_frame
body2.text = "This report summarizes the status of all legacy infrastructure components across regional data centers."
p2 = body2.add_paragraph()
p2.text = "Key focus areas include server decommissioning timelines, migration progress, and cost projections."
p2.space_before = Pt(12)
p3 = body2.add_paragraph()
p3.text = "All data reflects records as of September 30, 2024."
p3.space_before = Pt(12)

# ── Slide 3: Migration Timeline (Title + Content) ──
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
slide3.shapes.title.text = "Migration Timeline"
body3 = slide3.placeholders[1].text_frame
body3.text = "Phase 1 (Jan-Mar): Assessment and inventory of legacy servers"
for item in [
    "Phase 2 (Apr-Jun): Pilot migration of non-critical workloads",
    "Phase 3 (Jul-Sep): Production migration and validation",
    "Phase 4 (Oct-Dec): Decommissioning and final audit",
]:
    p = body3.add_paragraph()
    p.text = item
    p.space_before = Pt(8)

# ── Slide 4: Legacy Data Table (4 columns x 9 rows) ──
slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

# Add a title text box at top
title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Legacy Server Inventory"
p.alignment = PP_ALIGN.CENTER
run = p.runs[0]
run.font.size = Pt(24)
run.font.bold = True
run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

# Table data: header + 8 data rows = 9 rows total
# Rows 7-9 (index 6-8) are deprecated
table_data = [
    ["Server ID",    "Location",       "Status",       "Decommission Date"],
    ["SRV-1001",     "US-East-1",      "Active",       "2025-03-15"],
    ["SRV-1002",     "US-West-2",      "Active",       "2025-06-01"],
    ["SRV-1003",     "EU-Central-1",   "Migrating",    "2025-01-20"],
    ["SRV-1004",     "AP-Southeast-1", "Active",       "2025-04-30"],
    ["SRV-1005",     "US-East-2",      "Migrating",    "2025-02-28"],
    ["SRV-0087",     "US-West-1",      "Deprecated",   "2023-06-15"],
    ["SRV-0054",     "EU-West-1",      "Deprecated",   "2022-11-30"],
    ["SRV-0031",     "AP-Northeast-1", "Deprecated",   "2022-03-01"],
]

rows, cols = 9, 4
table_shape = slide4.shapes.add_table(rows, cols, Inches(0.5), Inches(1.0), Inches(9.0), Inches(5.0))
table = table_shape.table

# Set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(2.5)
table.columns[2].width = Inches(2.0)
table.columns[3].width = Inches(2.5)

# Populate table
for r_idx, row_data in enumerate(table_data):
    for c_idx, cell_text in enumerate(row_data):
        cell = table.cell(r_idx, c_idx)
        cell.text = cell_text
        para = cell.text_frame.paragraphs[0]
        for run in para.runs:
            run.font.size = Pt(11)
            if r_idx == 0:
                run.font.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            elif r_idx >= 6:
                # Deprecated rows in muted gray
                run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
                run.font.italic = True

# Style header row background
for c_idx in range(cols):
    cell = table.cell(0, c_idx)
    fill = cell.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

# ── Slide 5: Summary & Next Steps ──
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
slide5.shapes.title.text = "Summary & Next Steps"
body5 = slide5.placeholders[1].text_frame
body5.text = "5 active servers scheduled for migration by Q2 2025"
for item in [
    "3 deprecated servers pending removal from inventory records",
    "Budget allocation for Phase 4 decommissioning approved",
    "Next review scheduled for January 15, 2025",
]:
    p = body5.add_paragraph()
    p.text = item
    p.space_before = Pt(8)

# Save the presentation
prs.save(OUTPUT_PATH)
print(f"Presentation saved to {OUTPUT_PATH}")

# Launch LibreOffice Impress
def launch_gui(command: str, delay_sec: float = 1.0):
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)

launch_gui(f'libreoffice --impress "{OUTPUT_PATH}"', delay_sec=2.0)
print("LibreOffice Impress launched.")
