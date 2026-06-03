"""
Reward script for impress_tct_022:
Remove the bottom 3 rows from the table on slide 4 that contain deprecated data.

Scoring rubric (total 1.0):
  - 0.6: Table on slide 4 has exactly 6 rows (was 9, bottom 3 deleted)
  - 0.4: First 6 rows data intact after deletion
         (only awarded if row count changed from original 9)
"""

import pyautogui
import time

# Persistence hook: save any unsaved GUI edits
pyautogui.hotkey("ctrl", "s")
time.sleep(0.8)


def evaluate():
    score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        print("python-pptx not available")
        print("REWARD: 0.0")
        return

    pptx_path = "/home/user/Legacy_Data.pptx"

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"Cannot open presentation: {e}")
        print("REWARD: 0.0")
        return

    # Check slide count
    num_slides = len(prs.slides)
    if num_slides < 4:
        print(f"Only {num_slides} slides, need at least 4")
        print("REWARD: 0.0")
        return

    # Find table on slide 4 (index 3)
    slide4 = prs.slides[3]
    table = None
    for shape in slide4.shapes:
        if shape.has_table:
            table = shape.table
            break

    if table is None:
        print("No table found on slide 4")
        print("REWARD: 0.0")
        return

    num_rows = len(table.rows)
    num_cols = len(table.columns)
    print(f"Table dimensions: {num_rows} rows x {num_cols} cols")

    # Expected data for first 6 rows (row 0 = header, rows 1-5 = data)
    expected_rows = [
        ["Server ID", "Location", "Status", "Decommission Date"],
        ["SRV-1001", "US-East-1", "Active", "2025-03-15"],
        ["SRV-1002", "US-West-2", "Active", "2025-06-01"],
        ["SRV-1003", "EU-Central-1", "Migrating", "2025-01-20"],
        ["SRV-1004", "AP-Southeast-1", "Active", "2025-04-30"],
        ["SRV-1005", "US-East-2", "Migrating", "2025-02-28"],
    ]

    # Component 1 (0.6): Row count is exactly 6
    ORIGINAL_ROW_COUNT = 9
    if num_rows == 6:
        score += 0.6
        print("Row count check: PASS (6 rows)")
    else:
        print(f"Row count check: FAIL (expected 6, got {num_rows})")

    # Component 2 (0.4): First 6 rows data intact
    # Only awarded if rows were actually removed (row count != original 9)
    if num_rows != ORIGINAL_ROW_COUNT:
        rows_to_check = min(num_rows, 6)
        matching_rows = 0
        for i in range(rows_to_check):
            actual = [table.cell(i, c).text.strip() for c in range(min(num_cols, 4))]
            if i < len(expected_rows) and actual == expected_rows[i]:
                matching_rows += 1
            else:
                print(f"Row {i} mismatch: expected {expected_rows[i] if i < len(expected_rows) else 'N/A'}, got {actual}")

        if rows_to_check > 0:
            data_score = (matching_rows / 6) * 0.4
            score += data_score
            print(f"Data integrity: {matching_rows}/6 rows match (score +{data_score:.2f})")
        else:
            print("Data integrity: no rows to check")
    else:
        print("Data integrity: SKIPPED (row count unchanged from original 9)")

    score = round(min(score, 1.0), 2)
    print(f"REWARD: {score}")


if __name__ == "__main__":
    evaluate()
