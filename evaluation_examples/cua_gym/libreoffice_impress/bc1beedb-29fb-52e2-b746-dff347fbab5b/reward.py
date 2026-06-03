"""
FINAL REWARD SCRIPT - SUCCESS
Task: I'm reworking my presentation to make our logo pop a bit more. How can I turn the logo on the first slide a perfect 30 degrees clockwise?
Generated: 2025-08-07 09:16:36
Status: success
Model: o4-mini
Total Steps: 2
"""

import os
from pptx import Presentation

def verify_rotation(file_path):
    """
    Verifies that the logo (picture) on the first slide is rotated exactly 30 degrees clockwise.
    Progressive scoring:
      - 0.2 points: file exists
      - 0.1 points: at least one shape on first slide
      - 0.2 points: picture shape found
      - 0.5 points: picture rotation exactly 30 degrees
    Returns a float between 0.0 and 1.0 and prints detailed verification steps.
    """
    total_score = 0.0
    max_score = 1.0

    # 1. Check file existence
    print("Checking file existence...")
    if os.path.exists(file_path):
        print("✓ File exists (0.2 points)")
        total_score += 0.2
    else:
        print("✗ File not found (0 points)")
        print(f"REWARD: {min(total_score, max_score)}")
        return min(total_score, max_score)

    # 2. Load presentation
    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded with {len(prs.slides)} slide(s) (load success)")
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        print(f"REWARD: {min(total_score, max_score)}")
        return min(total_score, max_score)

    # 3. Verify shapes on first slide
    slide = prs.slides[0]
    print(f"Inspecting first slide: found {len(slide.shapes)} shape(s)")
    if len(slide.shapes) > 0:
        print("✓ Shapes present on first slide (0.1 points)")
        total_score += 0.1
    else:
        print("✗ No shapes found on first slide (0 points)")

    # 4. Find picture shape
    pic_shape = None
    for shape in slide.shapes:
        if hasattr(shape, 'shape_type') and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            pic_shape = shape
            break

    if pic_shape:
        print("✓ Picture shape found (0.2 points)")
        total_score += 0.2
        rotation = getattr(pic_shape, 'rotation', None)
        print(f"Picture rotation: {rotation}")
        # 5. Check rotation
        if rotation is not None and abs(rotation - 30.0) < 1e-6:
            print("✓ Rotation is exactly 30 degrees (0.5 points)")
            total_score += 0.5
        else:
            print("✗ Rotation is not 30 degrees (0 points)")
    else:
        print("✗ No picture shape found (0 points)")

    # 6. Final scoring
    final_score = min(total_score, max_score)
    print(f"Total score breakdown: {total_score}/{max_score}")
    print(f"Final score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score

if __name__ == '__main__':
    # Path to the given presentation file
    file_path = '/home/user/im_reworking_my_presentation_to_make_our_logo_pop_a_bit_more_how_can_i_turn_the_logo_on_the_first_sl.pptx'
    verify_rotation(file_path)

