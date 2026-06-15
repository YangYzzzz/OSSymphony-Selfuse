"""
Initial Setup: Create a 7-slide presentation with slide 3 titled 'Project Schedule' and empty content.
Task ID: impress_gf2_019
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_019'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Project Timeline"
    slide1.placeholders[1].text = "Q1 2025 Software Development Roadmap"

    # --- Slide 2: Project Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Project Overview"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "The Horizon Platform is a next-generation customer analytics solution designed to consolidate data streams from multiple touchpoints."
    p2 = tf2.add_paragraph()
    p2.text = "Key objectives include real-time dashboard reporting, automated anomaly detection, and seamless integration with existing CRM tools."
    p3 = tf2.add_paragraph()
    p3.text = "The project is led by Sarah Chen (Engineering) with cross-functional support from Marketing and Product teams."

    # --- Slide 3: Project Schedule (empty - task target) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf3 = txBox.text_frame
    p = tf3.paragraphs[0]
    p.text = "Project Schedule"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    # No table - this is where the agent should add the table

    # --- Slide 4: Team Members ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Team Members"
    tf4 = slide4.placeholders[1].text_frame
    members = [
        "Sarah Chen - Project Lead & Engineering Manager",
        "Marcus Johnson - Senior Backend Developer",
        "Priya Patel - Frontend Developer & UX Lead",
        "David Kim - QA Engineer & Test Automation",
        "Elena Rodriguez - DevOps & Infrastructure",
        "James Wright - Product Manager",
    ]
    tf4.text = members[0]
    for m in members[1:]:
        p = tf4.add_paragraph()
        p.text = m

    # --- Slide 5: Budget Summary ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Budget Summary"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Total Allocated Budget: $285,000"
    items = [
        "Personnel Costs: $180,000 (63%)",
        "Infrastructure & Cloud Services: $45,000 (16%)",
        "Software Licenses & Tools: $28,000 (10%)",
        "Training & Professional Development: $12,000 (4%)",
        "Contingency Reserve: $20,000 (7%)",
    ]
    for item in items:
        p = tf5.add_paragraph()
        p.text = item

    # --- Slide 6: Risks & Mitigations ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Risks & Mitigations"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Risk 1: Scope creep from stakeholder feature requests"
    p = tf6.add_paragraph()
    p.text = "Mitigation: Strict change control process with weekly scope reviews"
    p = tf6.add_paragraph()
    p.text = ""
    p = tf6.add_paragraph()
    p.text = "Risk 2: Integration complexity with legacy CRM systems"
    p = tf6.add_paragraph()
    p.text = "Mitigation: Early proof-of-concept testing and dedicated integration sprint"

    # --- Slide 7: Next Steps ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Next Steps"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "Complete Requirements phase by January 31, 2025"
    steps = [
        "Finalize system architecture and design documentation",
        "Set up CI/CD pipeline and staging environment",
        "Begin Sprint 1 development with core API endpoints",
        "Schedule bi-weekly stakeholder demos starting February 15",
    ]
    for s in steps:
        p = tf7.add_paragraph()
        p.text = s

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
