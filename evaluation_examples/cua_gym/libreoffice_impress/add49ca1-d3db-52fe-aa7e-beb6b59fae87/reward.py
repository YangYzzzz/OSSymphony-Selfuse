"""
Reward Script: Insert formatted project timeline table on slide 3
Task ID: impress_gf2_019
Domain: libreoffice_impress
Scoring:
  Component 1: Table exists on slide 3 with correct dimensions (7 rows x 5 cols) — 0.20
  Component 2: Header row has correct column names — 0.15
  Component 3: Data rows contain correct task names and status values — 0.25
  Component 4: Header row dark background with white bold text — 0.15
  Component 5: Completed rows (Requirements, Design) have green #D1FAE5 background — 0.15
  Component 6: Development row has yellow #FEF3C7, upcoming rows have white #FFFFFF — 0.10
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_019'


def persist_app_state(domain):
    """Save any unsaved LibreOffice changes before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_cell_fill(cell):
    """Extract solid fill color hex from a table cell, or None."""
    from pptx.oxml.ns import qn
    tcPr = cell._tc.tcPr
    if tcPr is not None:
        solidFill = tcPr.find(qn('a:solidFill'))
        if solidFill is not None:
            srgb = solidFill.find(qn('a:srgbClr'))
            if srgb is not None:
                return srgb.get('val').upper()
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]

    # Find table on slide 3
    table = None
    for shape in slide3.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            break

    # Component 1: Table exists on slide 3 with correct dimensions (0.20 points)
    try:
        if table is None:
            print("FAIL: Component 1 — No table found on slide 3")
        else:
            num_rows = len(table.rows)
            num_cols = len(table.columns)
            if num_rows == 7 and num_cols == 5:
                print(f"PASS: Component 1 — Table found: {num_rows}x{num_cols} (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 1 — Table is {num_rows}x{num_cols}, expected 7x5")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    if table is None:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Header row has correct column names (0.15 points)
    try:
        expected_headers = ['Task', 'Start Date', 'End Date', 'Duration', 'Status']
        actual_headers = [table.cell(0, c).text.strip() for c in range(min(len(table.columns), 5))]
        matches = sum(1 for a, e in zip(actual_headers, expected_headers) if a.lower() == e.lower())
        if matches == 5:
            print(f"PASS: Component 2 — All 5 headers correct: {actual_headers} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 — Headers: {actual_headers}, expected: {expected_headers} ({matches}/5 match)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Data rows contain correct task names and status values (0.25 points)
    try:
        expected_tasks = [
            ('Requirements', 'Completed'),
            ('Design', 'Completed'),
            ('Development', 'In Progress'),
            ('Testing', 'Upcoming'),
            ('Deployment', 'Upcoming'),
            ('Review', 'Upcoming'),
        ]
        task_matches = 0
        status_matches = 0
        for i, (exp_task, exp_status) in enumerate(expected_tasks):
            row_idx = i + 1  # skip header
            if row_idx >= len(table.rows):
                break
            actual_task = table.cell(row_idx, 0).text.strip()
            actual_status = table.cell(row_idx, 4).text.strip()
            if actual_task.lower() == exp_task.lower():
                task_matches += 1
            else:
                print(f"  Row {row_idx}: task={actual_task!r}, expected={exp_task!r}")
            if actual_status.lower() == exp_status.lower():
                status_matches += 1
            else:
                print(f"  Row {row_idx}: status={actual_status!r}, expected={exp_status!r}")

        # Award partial: tasks worth 0.15, statuses worth 0.10
        task_score = (task_matches / 6) * 0.15
        status_score = (status_matches / 6) * 0.10
        comp3_score = task_score + status_score
        if comp3_score > 0:
            print(f"PASS: Component 3 — Tasks {task_matches}/6, Statuses {status_matches}/6 ({comp3_score:.2f} pts)")
            total_score += comp3_score
        else:
            print(f"FAIL: Component 3 — No task/status matches")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Header row has dark background with white bold text (0.15 points)
    try:
        header_fill = get_cell_fill(table.cell(0, 0))
        # Check for a dark fill (not white, not light)
        header_is_dark = False
        if header_fill is not None:
            # Parse hex and check luminance
            r_val = int(header_fill[0:2], 16)
            g_val = int(header_fill[2:4], 16)
            b_val = int(header_fill[4:6], 16)
            luminance = 0.299 * r_val + 0.587 * g_val + 0.114 * b_val
            header_is_dark = luminance < 128

        # Check for bold white text in header
        header_bold = False
        header_white = False
        for c in range(min(len(table.columns), 5)):
            cell = table.cell(0, c)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    if run.font.bold is True:
                        header_bold = True
                    try:
                        if run.font.color.type is not None:
                            rgb_str = str(run.font.color.rgb).upper()
                            if rgb_str == 'FFFFFF':
                                header_white = True
                    except:
                        pass

        comp4_score = 0.0
        if header_is_dark:
            comp4_score += 0.075
        if header_bold and header_white:
            comp4_score += 0.075

        if comp4_score > 0:
            print(f"PASS: Component 4 — Header fill={header_fill}, dark={header_is_dark}, bold={header_bold}, white_text={header_white} ({comp4_score:.3f} pts)")
            total_score += comp4_score
        else:
            print(f"FAIL: Component 4 — Header fill={header_fill}, dark={header_is_dark}, bold={header_bold}, white_text={header_white}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Completed rows (Requirements row 1, Design row 2) have green #D1FAE5 background (0.15 points)
    try:
        green_count = 0
        for row_idx in [1, 2]:  # Requirements and Design
            fill = get_cell_fill(table.cell(row_idx, 0))
            if fill is not None and fill == 'D1FAE5':
                green_count += 1
                print(f"  Row {row_idx} ({table.cell(row_idx, 0).text}): fill={fill} ✓")
            else:
                print(f"  Row {row_idx} ({table.cell(row_idx, 0).text}): fill={fill}, expected D1FAE5")

        if green_count == 2:
            print(f"PASS: Component 5 — Both completed rows have green background (0.15 pts)")
            total_score += 0.15
        elif green_count == 1:
            print(f"PARTIAL: Component 5 — 1/2 completed rows have green background (0.075 pts)")
            total_score += 0.075
        else:
            print(f"FAIL: Component 5 — No completed rows have green background")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Development row yellow #FEF3C7, upcoming rows white #FFFFFF (0.10 points)
    try:
        comp6_score = 0.0

        # Development row (row 3) should be yellow
        dev_fill = get_cell_fill(table.cell(3, 0))
        if dev_fill is not None and dev_fill == 'FEF3C7':
            comp6_score += 0.05
            print(f"  Development row: fill={dev_fill} ✓")
        else:
            print(f"  Development row: fill={dev_fill}, expected FEF3C7")

        # Upcoming rows (4, 5, 6) should be white
        white_count = 0
        for row_idx in [4, 5, 6]:
            fill = get_cell_fill(table.cell(row_idx, 0))
            if fill is not None and fill == 'FFFFFF':
                white_count += 1
            else:
                print(f"  Row {row_idx} ({table.cell(row_idx, 0).text}): fill={fill}, expected FFFFFF")

        if white_count == 3:
            comp6_score += 0.05
            print(f"  All 3 upcoming rows white ✓")

        if comp6_score > 0:
            print(f"PASS: Component 6 — Conditional formatting ({comp6_score:.2f} pts)")
            total_score += comp6_score
        else:
            print(f"FAIL: Component 6 — No conditional formatting matches")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    final_score = min(round(total_score, 4), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state('libreoffice_impress')

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
