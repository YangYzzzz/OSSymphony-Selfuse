"""
Initial Setup: Product Roadmap Presentation (9 slides, no transitions)
Task ID: impress_tm_022
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_022'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_text(slide, title_text, subtitle_text=None):
    """Set title and optional subtitle on a slide."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    if subtitle_text and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle_text


def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=None):
    """Add a textbox with formatted text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=12):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(font_size)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    add_title_text(slide1, "Nextera Product Roadmap 2025",
                   "Strategic Vision & Quarterly Milestones")

    # --- Slide 2: Q1 Objectives ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide2, "Q1 2025 — Foundation & Core Platform")
    add_bullet_list(slide2, 0.8, 2.0, 8.5, 4.5, [
        "Launch redesigned authentication system with OAuth 2.1 support",
        "Complete migration from PostgreSQL 14 to PostgreSQL 16",
        "Deploy real-time analytics dashboard for enterprise clients",
        "Reduce average API response time to under 120ms",
        "Onboard 15 beta partners for the developer portal",
    ])

    # --- Slide 3: Q1 Key Metrics ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_textbox(slide3, 1.0, 0.5, 8.0, 1.0, "Q1 Key Performance Targets",
                font_size=28, bold=True, color=RGBColor(0x1F, 0x49, 0x7D))
    # Add table
    table_shape = slide3.shapes.add_table(5, 3, Inches(1.0), Inches(1.8), Inches(8.0), Inches(3.5))
    table = table_shape.table
    headers = ["Metric", "Current", "Q1 Target"]
    data = [
        ["Monthly Active Users", "142,000", "185,000"],
        ["API Uptime", "99.92%", "99.97%"],
        ["Avg Response Time", "187ms", "120ms"],
        ["Enterprise Clients", "38", "52"],
    ]
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
        for run in table.cell(0, c).text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row in enumerate(data, 1):
        for c, val in enumerate(row):
            table.cell(r, c).text = val

    # --- Slide 4: Q2 Strategy ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide4, "Q2 2025 — Growth & Integrations")
    add_bullet_list(slide4, 0.8, 2.0, 8.5, 4.5, [
        "Release Salesforce CRM bi-directional sync module",
        "Implement AI-powered search across knowledge base (RAG pipeline)",
        "Launch mobile SDK for iOS and Android platforms",
        "Introduce tiered pricing: Starter ($29/mo), Pro ($99/mo), Enterprise (custom)",
        "Expand engineering team by 8 full-stack developers",
    ])

    # --- Slide 5: Q2 Architecture ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_textbox(slide5, 1.0, 0.5, 8.0, 1.0, "Q2 Architecture Evolution",
                font_size=28, bold=True, color=RGBColor(0x1F, 0x49, 0x7D))
    add_bullet_list(slide5, 1.0, 1.8, 8.0, 5.0, [
        "Microservices decomposition: monolith split into 12 bounded contexts",
        "Event-driven architecture with Apache Kafka for inter-service messaging",
        "Redis cluster for session management and caching (projected 40% latency reduction)",
        "Kubernetes autoscaling policies refined based on Q1 traffic analysis",
        "GraphQL gateway replacing legacy REST endpoints for frontend clients",
        "Terraform modules standardized for multi-region deployment (US-East, EU-West, AP-Southeast)",
    ], font_size=11)

    # --- Slide 6: Q2 Revenue Projections ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_textbox(slide6, 1.0, 0.5, 8.0, 1.0, "Q2 Revenue Projections",
                font_size=28, bold=True, color=RGBColor(0x1F, 0x49, 0x7D))
    table_shape2 = slide6.shapes.add_table(5, 3, Inches(1.0), Inches(1.8), Inches(8.0), Inches(3.5))
    table2 = table_shape2.table
    headers2 = ["Revenue Stream", "Q1 Actual", "Q2 Forecast"]
    data2 = [
        ["SaaS Subscriptions", "$1,240,000", "$1,580,000"],
        ["Enterprise Licensing", "$890,000", "$1,120,000"],
        ["Professional Services", "$340,000", "$410,000"],
        ["API Usage Fees", "$175,000", "$260,000"],
    ]
    for c, h in enumerate(headers2):
        table2.cell(0, c).text = h
        for run in table2.cell(0, c).text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row in enumerate(data2, 1):
        for c, val in enumerate(row):
            table2.cell(r, c).text = val

    # --- Slide 7: Q3 Innovation ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide7, "Q3 2025 — AI Features & Scale")
    add_bullet_list(slide7, 0.8, 2.0, 8.5, 4.5, [
        "Launch AI-generated report builder with natural language queries",
        "Implement predictive analytics for customer churn (target: 85% accuracy)",
        "Deploy automated compliance scanning for SOC2 and GDPR workflows",
        "Scale infrastructure to handle 500K concurrent users",
        "Introduce white-label solution for reseller partners",
    ])

    # --- Slide 8: Q4 Vision ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide8, "Q4 2025 — Market Expansion & Polish")
    add_bullet_list(slide8, 0.8, 2.0, 8.5, 4.5, [
        "Enter APAC market with localized versions (Japanese, Korean, Mandarin)",
        "Launch marketplace for third-party integrations and plugins",
        "Achieve ISO 27001 certification for enterprise security requirements",
        "Release desktop application for Windows and macOS (Electron-based)",
        "Target $15M ARR by end of Q4 2025",
    ])

    # --- Slide 9: Summary & Next Steps ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_textbox(slide9, 1.0, 0.5, 8.0, 1.0, "Summary & Next Steps",
                font_size=28, bold=True, color=RGBColor(0x1F, 0x49, 0x7D))
    add_bullet_list(slide9, 1.0, 1.8, 8.0, 5.0, [
        "Q1: Solidify platform reliability and developer experience",
        "Q2: Accelerate growth through integrations and mobile presence",
        "Q3: Differentiate with AI-powered features at enterprise scale",
        "Q4: Expand globally and build partner ecosystem",
        "",
        "Key risks: Engineering hiring pipeline, Kafka migration complexity",
        "Budget allocation: 45% Engineering, 25% Sales, 15% Marketing, 15% Operations",
        "Executive review scheduled: March 28, 2025",
    ], font_size=11)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
