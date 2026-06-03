"""
Initial Setup: Geology Lecture presentation with 10 slides, no transitions
Task ID: impress_teach_017
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_017'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.paragraphs[0].text = point
        else:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(prs,
                    "Introduction to Physical Geology",
                    "GEOL 101 — Professor Maria Alvarez\nFall Semester 2025")

    # Slide 2: Course Overview
    add_content_slide(prs, "Course Overview", [
        "Understanding Earth's dynamic processes and materials",
        "Weekly field observations and laboratory exercises",
        "Topics span from minerals to plate tectonics",
        "Assessment: 2 exams, 1 field report, weekly quizzes",
        "Office hours: Tuesdays & Thursdays 2:00–3:30 PM",
    ])

    # Slide 3: Earth's Structure
    add_content_slide(prs, "Earth's Internal Structure", [
        "Inner core: solid iron-nickel alloy (5,150–6,371 km depth)",
        "Outer core: liquid iron-nickel (2,890–5,150 km depth)",
        "Mantle: silicate rock undergoing convection (35–2,890 km)",
        "Crust: oceanic (5–10 km) vs. continental (30–70 km)",
        "Lithosphere and asthenosphere boundary controls tectonics",
    ])

    # Slide 4: Minerals and Crystallography
    add_content_slide(prs, "Minerals and Crystallography", [
        "A mineral: naturally occurring, inorganic, ordered atomic structure",
        "Mohs hardness scale: talc (1) to diamond (10)",
        "Silicate minerals make up >90% of Earth's crust",
        "Key groups: feldspars, quartz, micas, olivine, pyroxenes",
        "Crystal systems: cubic, tetragonal, hexagonal, orthorhombic, monoclinic, triclinic",
    ])

    # Slide 5: Igneous Rocks
    add_content_slide(prs, "Igneous Rocks", [
        "Formed from cooling and solidification of magma or lava",
        "Intrusive (plutonic): slow cooling → coarse-grained (granite, gabbro)",
        "Extrusive (volcanic): rapid cooling → fine-grained (basalt, rhyolite)",
        "Bowen's Reaction Series predicts crystallization order",
        "Texture classifications: phaneritic, aphanitic, porphyritic, glassy",
    ])

    # Slide 6: Sedimentary Rocks
    add_content_slide(prs, "Sedimentary Rocks", [
        "Formed by deposition, compaction, and cementation of sediments",
        "Clastic: sandstone, shale, conglomerate, siltstone",
        "Chemical: limestone, evaporites (halite, gypsum)",
        "Organic: coal, some limestones (coquina, chalk)",
        "Strata and fossils provide crucial geological history records",
    ])

    # Slide 7: Metamorphic Rocks
    add_content_slide(prs, "Metamorphic Rocks", [
        "Produced by heat, pressure, or chemically active fluids",
        "Foliated: slate, phyllite, schist, gneiss (increasing grade)",
        "Non-foliated: marble (from limestone), quartzite (from sandstone)",
        "Contact metamorphism: localized heat near igneous intrusions",
        "Regional metamorphism: large-scale tectonic compression zones",
    ])

    # Slide 8: Plate Tectonics
    add_content_slide(prs, "Plate Tectonics", [
        "Lithosphere divided into ~15 major tectonic plates",
        "Divergent boundaries: mid-ocean ridges, rift valleys (East Africa)",
        "Convergent boundaries: subduction zones, mountain building (Himalayas)",
        "Transform boundaries: lateral sliding (San Andreas Fault)",
        "Driving forces: mantle convection, ridge push, slab pull",
    ])

    # Slide 9: Earthquakes and Seismology
    add_content_slide(prs, "Earthquakes and Seismology", [
        "Caused by sudden release of energy along faults",
        "Focus (hypocenter) vs. epicenter: depth classification matters",
        "Seismic waves: P-waves (compressional), S-waves (shear), surface waves",
        "Moment magnitude scale replaced older Richter scale",
        "Seismic hazard mapping guides building codes worldwide",
    ])

    # Slide 10: Volcanism and Volcanic Hazards
    add_content_slide(prs, "Volcanism and Volcanic Hazards", [
        "Shield volcanoes: broad, gentle slopes (Mauna Loa, Kīlauea)",
        "Stratovolcanoes: steep composite cones (Mt. St. Helens, Fuji)",
        "Cinder cones: small, steep, short-lived eruptions (Parícutin)",
        "Hazards: pyroclastic flows, lahars, tephra fallout, gas emissions",
        "Monitoring: seismicity, ground deformation, gas chemistry, thermal imaging",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
