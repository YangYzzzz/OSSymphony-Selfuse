"""
Initial Setup: Create a 5-slide Yearly Comparison presentation with a 2D column chart on slide 3.
Task ID: impress_tct_049
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_049'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Yearly Performance Comparison"
    slide1.placeholders[1].text = "Revenue Analysis: 2024 vs 2025"

    # ---- Slide 2: Overview ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Key Highlights"
    points = [
        "Total revenue grew by 12% year-over-year across all quarters",
        "Q3 showed the strongest growth driven by the new product launch",
        "Q1 remained stable despite seasonal market slowdown",
        "Q4 exceeded projected targets by 8%, closing the year on a high note",
    ]
    for pt in points:
        p = body2.add_paragraph()
        p.text = pt
        p.level = 1

    # ---- Slide 3: 2D Column Chart ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title text box
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Quarterly Revenue Comparison"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Chart data: 4 categories, 2 series
    chart_data = CategoryChartData()
    chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
    chart_data.add_series('Revenue 2024', (245000, 312000, 287000, 356000))
    chart_data.add_series('Revenue 2025', (268000, 341000, 378000, 402000))

    # Add 2D clustered column chart
    chart_frame = slide3.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.0), Inches(1.2),
        Inches(11.0), Inches(5.8),
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # Style the chart series
    series0 = chart.series[0]
    series0.format.fill.solid()
    series0.format.fill.fore_color.rgb = RGBColor(0x4472, 0xC4, 0x00)[:3] if False else RGBColor(0x44, 0x72, 0xC4)
    series1 = chart.series[1]
    series1.format.fill.solid()
    series1.format.fill.fore_color.rgb = RGBColor(0xED, 0x7D, 0x31)

    # ---- Slide 4: Analysis ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Quarterly Analysis"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Performance Breakdown"
    details = [
        "Q1: Steady start with 9.4% growth driven by enterprise contracts",
        "Q2: Strong momentum with 9.3% increase from expanded retail partnerships",
        "Q3: Peak performance at 31.7% growth fueled by product launch",
        "Q4: Solid close with 12.9% growth from holiday season demand",
    ]
    for d in details:
        p = body4.add_paragraph()
        p.text = d
        p.level = 1

    # ---- Slide 5: Summary ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Summary & Outlook"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Looking Ahead to 2026"
    conclusions = [
        "Overall revenue increased from $1.2M to $1.389M, a 15.8% annual gain",
        "Growth trajectory suggests continued upward trend for 2026",
        "Investment in Q3 product launches proved highly effective",
        "Recommendation: Increase Q1 marketing spend to reduce seasonal dip",
    ]
    for c in conclusions:
        p = body5.add_paragraph()
        p.text = c
        p.level = 1

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
