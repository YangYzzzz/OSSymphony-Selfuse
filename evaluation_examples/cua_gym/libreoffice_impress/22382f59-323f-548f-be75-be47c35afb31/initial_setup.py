"""
Initial Setup: Create a 7-slide Regional Sales presentation with slide 4 having
title 'Regional Performance' and an empty content area (no table).
Task ID: impress_gf3_015
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_015'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Regional Sales Report"
    slide1.placeholders[1].text = "FY 2025 Annual Review"

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Overall revenue grew 12% year-over-year across all regions."
    p2 = tf2.add_paragraph()
    p2.text = "North region led growth with a 15% increase in Q4."
    p3 = tf2.add_paragraph()
    p3.text = "South and East regions showed steady improvement throughout the year."

    # --- Slide 3: Market Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Market Overview"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "The domestic market expanded significantly in 2025."
    p3a = tf3.add_paragraph()
    p3a.text = "Consumer spending increased 8% driven by favorable economic conditions."
    p3b = tf3.add_paragraph()
    p3b.text = "Key competitors saw slower growth, providing opportunity for market share gains."

    # --- Slide 4: Regional Performance (EMPTY - no table) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    tf4 = txBox.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Regional Performance"
    p4.alignment = PP_ALIGN.LEFT
    run4 = p4.runs[0]
    run4.font.size = Pt(36)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 5: Channel Analysis ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Channel Analysis"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Online channels contributed 45% of total revenue, up from 38% in 2024."
    p5a = tf5.add_paragraph()
    p5a.text = "Retail partnerships remained stable with 3% growth."
    p5b = tf5.add_paragraph()
    p5b.text = "B2B direct sales showed the strongest margin improvement at 22%."

    # --- Slide 6: Key Initiatives ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Key Initiatives for 2026"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Expand into West region with targeted marketing campaigns."
    p6a = tf6.add_paragraph()
    p6a.text = "Invest in supply chain optimization to reduce delivery times by 20%."
    p6b = tf6.add_paragraph()
    p6b.text = "Launch customer loyalty program in Q2 to improve retention rates."

    # --- Slide 7: Thank You ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[0])
    slide7.shapes.title.text = "Thank You"
    slide7.placeholders[1].text = "Questions & Discussion"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
