"""
Reward Script: Insert a 5x4 table on slide 4 with sales data and formatted header
Task ID: impress_gf3_015
Domain: libreoffice_impress
Scoring:
  Component 1: Table exists on slide 4 with correct dimensions (5 cols, 4 rows) — 0.20
  Component 2: Header row text matches expected labels — 0.20
  Component 3: Data rows contain correct values — 0.30
  Component 4: Header cells have dark green (#1A5C1A) background fill — 0.15
  Component 5: Header text is bold and white (#FFFFFF) — 0.15
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_015'

EXPECTED_HEADERS = ['Region', 'Q1 Sales', 'Q2 Sales', 'Q3 Sales', 'Q4 Sales']
EXPECTED_DATA = [
    ['North', '120', '135', '115', '148'],
    ['South', '98', '107', '122', '131'],
    ['East', '145', '152', '138', '165'],
]


def find_table_on_slide(slide):
    """Find the first table shape on a slide, or None."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return shape.table
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide4 = prs.slides[3]  # 0-indexed

    # Component 1: Table exists on slide 4 with correct dimensions (0.20 points)
    try:
        table = find_table_on_slide(slide4)
        if table is None:
            print("FAIL: Component 1 — No table found on slide 4")
        else:
            n_rows = len(table.rows)
            n_cols = len(table.columns)
            if n_rows == 4 and n_cols == 5:
                print(f"PASS: Component 1 — Table found with {n_rows} rows x {n_cols} cols (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 1 — Table dimensions {n_rows}x{n_cols}, expected 4x5")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Gate: if no table, cannot check further components
    if table is None:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: Header row text content (0.20 points)
    try:
        actual_headers = [table.cell(0, c).text.strip() for c in range(len(table.columns))]
        matching = sum(1 for a, e in zip(actual_headers, EXPECTED_HEADERS) if a == e)
        if matching == 5:
            print(f"PASS: Component 2 — All 5 header cells match: {actual_headers} (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 2 — {matching}/5 headers match. Actual: {actual_headers}, Expected: {EXPECTED_HEADERS}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Data rows contain correct values (0.30 points)
    try:
        correct_cells = 0
        total_cells = 15  # 3 rows x 5 cols
        for r_idx, expected_row in enumerate(EXPECTED_DATA):
            actual_row = [table.cell(r_idx + 1, c).text.strip() for c in range(len(table.columns))]
            for c_idx in range(5):
                if actual_row[c_idx] == expected_row[c_idx]:
                    correct_cells += 1
                else:
                    print(f"  MISMATCH: Row {r_idx+1} Col {c_idx}: expected '{expected_row[c_idx]}', got '{actual_row[c_idx]}'")

        if correct_cells == total_cells:
            print(f"PASS: Component 3 — All {total_cells} data cells correct (0.30 pts)")
            total_score += 0.30
        elif correct_cells >= 10:
            partial = round(0.30 * (correct_cells / total_cells), 2)
            print(f"PARTIAL: Component 3 — {correct_cells}/{total_cells} data cells correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — Only {correct_cells}/{total_cells} data cells correct")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Header cells have dark green (#1A5C1A) background fill (0.15 points)
    try:
        green_count = 0
        for c in range(5):
            cell = table.cell(0, c)
            fill = cell.fill
            if fill.type is not None:
                try:
                    rgb = str(fill.fore_color.rgb).upper()
                    if rgb == '1A5C1A':
                        green_count += 1
                    else:
                        print(f"  Header cell (0,{c}): fill color {rgb}, expected 1A5C1A")
                except Exception:
                    print(f"  Header cell (0,{c}): could not read fill color")
            else:
                print(f"  Header cell (0,{c}): no fill set")

        if green_count == 5:
            print(f"PASS: Component 4 — All 5 header cells have #1A5C1A fill (0.15 pts)")
            total_score += 0.15
        elif green_count >= 3:
            partial = round(0.15 * (green_count / 5), 2)
            print(f"PARTIAL: Component 4 — {green_count}/5 header cells with correct fill ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — Only {green_count}/5 header cells have correct fill")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Header text is bold and white (#FFFFFF) (0.15 points)
    try:
        bold_white_count = 0
        for c in range(5):
            cell = table.cell(0, c)
            is_bold = False
            is_white = False
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.bold is True:
                        is_bold = True
                    try:
                        if run.font.color.type is not None:
                            rgb = str(run.font.color.rgb).upper()
                            if rgb == 'FFFFFF':
                                is_white = True
                    except Exception:
                        pass

            if is_bold and is_white:
                bold_white_count += 1
            else:
                print(f"  Header cell (0,{c}): bold={is_bold}, white={is_white}")

        if bold_white_count == 5:
            print(f"PASS: Component 5 — All 5 header cells have bold white text (0.15 pts)")
            total_score += 0.15
        elif bold_white_count >= 3:
            partial = round(0.15 * (bold_white_count / 5), 2)
            print(f"PARTIAL: Component 5 — {bold_white_count}/5 header cells with bold white text ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 — Only {bold_white_count}/5 header cells have bold white text")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Main entrypoint
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
