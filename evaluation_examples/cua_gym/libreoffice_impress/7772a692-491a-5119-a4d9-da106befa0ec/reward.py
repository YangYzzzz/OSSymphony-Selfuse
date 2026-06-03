"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 16 feels a bit loud—please take the body text and switch its font colour to Gray 50% (#808080) and set it in Italic so it blends better with the rest of the deck.
Generated: 2025-09-10 14:39:04
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.dml.color import RGBColor


def verify_slide16_body_text_format(file_path: str) -> float:
    """Verify that slide 16 body text is italic and coloured Gray 50% (#808080).

    Progressive scoring (0‒1):
      • 0.5 points for correct colour proportion
      • 0.5 points for correct italic proportion
    """
    print(f"Verifying task on file: {file_path}")

    # ---------- Basic file checks (no points awarded) ----------
    if not os.path.exists(file_path):
        print("✗ File does not exist")
        return 0.0
    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        return 0.0

    if len(prs.slides) < 16:
        print("✗ Presentation has fewer than 16 slides – cannot verify slide 16")
        return 0.0

    # ---------- Locate slide 16 (index 15) ----------
    slide = prs.slides[15]
    print("✓ Found slide 16")

    # ---------- Gather body runs (exclude title placeholders) ----------
    body_runs = []
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        ph = getattr(shape, "placeholder_format", None)
        if ph and ph.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            # Skip title placeholders – only body text should be verified
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text and run.text.strip():
                    body_runs.append(run)

    total_runs = len(body_runs)
    print(f"Total body text runs found: {total_runs}")
    if total_runs == 0:
        print("✗ No body text runs found on slide 16")
        return 0.0

    # ---------- Verification for each run ----------
    target_rgb = RGBColor(0x80, 0x80, 0x80)  # Gray 50% (#808080)
    colour_ok_count = 0
    italic_ok_count = 0

    for idx, run in enumerate(body_runs, start=1):
        font = run.font

        # --- Colour check ---
        colour_ok = False
        if font.color is not None and font.color.type == 1 and font.color.rgb is not None:
            # type==1 -> RGB colour
            colour_ok = font.color.rgb == target_rgb
        # (Theme or other colour types are considered incorrect for this strict check)

        # --- Italic check ---
        italic_ok = font.italic is True

        # Count successes
        if colour_ok:
            colour_ok_count += 1
        if italic_ok:
            italic_ok_count += 1

        print(
            f"Run {idx}: text='{run.text[:30]}', colour_ok={colour_ok}, italic_ok={italic_ok}"
        )

    # ---------- Progressive scoring ----------
    colour_ratio = colour_ok_count / total_runs
    italic_ratio = italic_ok_count / total_runs

    print(
        f"Colour ratio: {colour_ratio:.0%} (correct {colour_ok_count}/{total_runs})"
    )
    print(
        f"Italic ratio: {italic_ratio:.0%} (correct {italic_ok_count}/{total_runs})"
    )

    score = 0.5 * colour_ratio + 0.5 * italic_ratio
    score = round(min(score, 1.0), 2)

    print(f"Score calculated: {score}")
    return score


if __name__ == "__main__":
    FILE_PATH = "/home/user/slide_16_feels_a_bit_loudplease_take_the_body_text_and_switch_its_font_colour_to_gray_50_808080_and__golden.pptx"
    reward = verify_slide16_body_text_format(FILE_PATH)
    print(f"REWARD: {reward}")
