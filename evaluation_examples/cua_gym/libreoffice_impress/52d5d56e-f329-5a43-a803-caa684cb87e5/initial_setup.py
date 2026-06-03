"""
Initial Setup: Create a 10-slide Photo Story presentation with default layouts
Task ID: impress_gf2_030
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_030'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard widescreen: 33.87cm x 19.05cm
    prs.slide_width = Cm(33.87)
    prs.slide_height = Cm(19.05)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "A Journey Through the Alps"
    slide1.placeholders[1].text = "Photography Collection by Elena Vasquez\nSpring 2025"

    # --- Slide 2: Title + Content ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "The Route Overview"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Our expedition covered five countries over 21 days"
    p = tf2.add_paragraph()
    p.text = "Starting in Chamonix, France and ending in Innsbruck, Austria"
    p = tf2.add_paragraph()
    p.text = "Total distance: 1,240 km through mountain passes"
    p = tf2.add_paragraph()
    p.text = "Highest elevation reached: 3,842m at Aiguille du Midi"

    # --- Slide 3: Title + Content ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Day 1-3: Chamonix Valley"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Mont Blanc towered above us at 4,808 meters"
    p = tf3.add_paragraph()
    p.text = "Morning fog rolling through the Arve River valley"
    p = tf3.add_paragraph()
    p.text = "Alpine wildflowers in full bloom along the hiking trails"
    p = tf3.add_paragraph()
    p.text = "Traditional Savoyard fondue dinner at Refuge du Montenvers"

    # --- Slide 4: Title + Content ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Day 4-6: The Matterhorn"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Zermatt offered the iconic view of the Matterhorn at sunrise"
    p = tf4.add_paragraph()
    p.text = "Golden hour light at 5:47 AM casting long shadows"
    p = tf4.add_paragraph()
    p.text = "Gornergrat railway took us to 3,089m for panoramic shots"

    # --- Slide 5: Title + Content ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Day 7-9: Swiss Lake District"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Lake Lucerne mirror reflections at dawn"
    p = tf5.add_paragraph()
    p.text = "Pilatus summit reached via the world's steepest cogwheel railway"
    p = tf5.add_paragraph()
    p.text = "Interlaken paragliding captured with 70-200mm telephoto lens"

    # --- Slide 6: Title + Content ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Day 10-12: Dolomites, Italy"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Tre Cime di Lavaredo at sunset — the most photographed peaks"
    p = tf6.add_paragraph()
    p.text = "Seceda ridgeline with dramatic cloud inversions below"
    p = tf6.add_paragraph()
    p.text = "Val di Funes church with autumn-colored larches"

    # --- Slide 7: Title + Content ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Day 13-15: Bavarian Alps"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "Neuschwanstein Castle framed by alpine peaks"
    p = tf7.add_paragraph()
    p.text = "Eibsee Lake crystal-clear waters with Zugspitze backdrop"
    p = tf7.add_paragraph()
    p.text = "Traditional Bavarian village of Mittenwald at golden hour"

    # --- Slide 8: Title + Content ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Day 16-18: Austrian Tyrol"
    tf8 = slide8.placeholders[1].text_frame
    tf8.text = "Innsbruck's colorful facades along the Inn River"
    p = tf8.add_paragraph()
    p.text = "Stubai Glacier ice cave formations in blue light"
    p = tf8.add_paragraph()
    p.text = "Nordkette cable car ascending directly from city center"

    # --- Slide 9: Title + Content ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Equipment & Techniques"
    tf9 = slide9.placeholders[1].text_frame
    tf9.text = "Camera: Sony A7R V with 24-70mm f/2.8 GM II"
    p = tf9.add_paragraph()
    p.text = "Drone: DJI Mavic 3 Pro for aerial landscape shots"
    p = tf9.add_paragraph()
    p.text = "Filters: NiSi 10-stop ND for long exposures of waterfalls"
    p = tf9.add_paragraph()
    p.text = "Post-processing: Adobe Lightroom with custom mountain preset"

    # --- Slide 10: Title + Content ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "Thank You"
    tf10 = slide10.placeholders[1].text_frame
    tf10.text = "All photos available at evasquez-photography.com/alps"
    p = tf10.add_paragraph()
    p.text = "Contact: elena@vasquez-photography.com"
    p = tf10.add_paragraph()
    p.text = "Instagram: @elena_mountain_lens"
    p = tf10.add_paragraph()
    p.text = "Next expedition: Patagonia, November 2025"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
