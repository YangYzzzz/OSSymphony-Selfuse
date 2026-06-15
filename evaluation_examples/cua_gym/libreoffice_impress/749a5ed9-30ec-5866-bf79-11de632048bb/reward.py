"""
Reward Script: Lab safety rules slide with red triangle and numbered list
Task ID: impress_teach_086
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Triangle/arrow shape exists on slide 2
  Component 2 (0.15): Triangle fill color is #F44336
  Component 3 (0.30): All 6 safety rules text present on slide 2
  Component 4 (0.15): Safety rules font size is ~20pt (254000 EMU)
  Component 5 (0.15): Triangle positioned above the rules text box
"""

import os
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_086'

EXPECTED_RULES = [
    "1. Wear safety goggles at all times",
    "2. No food or drink in the lab",
    "3. Report all spills immediately",
    "4. Know the location of fire extinguishers",
    "5. Never work alone in the lab",
    "6. Dispose of chemicals properly",
]


def find_triangle_shape(slide):
    """Find a triangle/warning shape on the slide (AUTO_SHAPE type)."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # Check if it's a triangle-like shape
            try:
                auto_type = shape.auto_shape_type
                # ISOSCELES_TRIANGLE = 7, but also accept other triangle types
                if auto_type is not None and 'TRIANGLE' in str(auto_type):
                    return shape
            except Exception:
                pass
            # Fallback: check shape name for triangle indicators
            name_lower = shape.name.lower()
            if 'triangle' in name_lower or 'warning' in name_lower:
                return shape
    return None


def find_rules_textbox(slide):
    """Find the text box containing the safety rules (not the title)."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            # Look for the textbox that contains safety rule content
            if "Wear safety goggles" in text or "No food or drink" in text:
                return shape
            # Also check if it has numbered items
            if text.count("\n") >= 4 and ("1." in text or "2." in text):
                return shape
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 2")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]  # 0-indexed, slide 2

    # Component 1: Triangle shape exists on slide 2 (0.25 points)
    try:
        triangle = find_triangle_shape(slide2)
        if triangle is not None:
            print(f"PASS: Component 1 — Triangle shape found: '{triangle.name}' (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — No triangle shape found on slide 2")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Triangle fill color is #F44336 (0.15 points)
    try:
        if triangle is not None:
            fill = triangle.fill
            if fill.type is not None and fill.type == 1:  # SOLID fill
                color_rgb = str(fill.fore_color.rgb).upper()
                if color_rgb == "F44336":
                    print(f"PASS: Component 2 — Triangle fill color is #F44336 (0.15 pts)")
                    total_score += 0.15
                else:
                    print(f"FAIL: Component 2 — Triangle fill color is #{color_rgb}, expected #F44336")
            else:
                print(f"FAIL: Component 2 — Triangle fill is not solid (type={fill.type})")
        else:
            print(f"FAIL: Component 2 — No triangle shape to check fill color")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: All 6 safety rules text present on slide 2 (0.30 points)
    try:
        rules_box = find_rules_textbox(slide2)
        if rules_box is not None:
            found_text = rules_box.text_frame.text
            rules_found = 0
            for rule in EXPECTED_RULES:
                if rule in found_text:
                    rules_found += 1
                else:
                    print(f"  MISS: Rule not found: '{rule}'")

            if rules_found == 6:
                print(f"PASS: Component 3 — All 6 safety rules found (0.30 pts)")
                total_score += 0.30
            elif rules_found >= 4:
                partial = round(0.30 * (rules_found / 6), 2)
                print(f"PARTIAL: Component 3 — {rules_found}/6 rules found ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 3 — Only {rules_found}/6 rules found")
        else:
            print(f"FAIL: Component 3 — No text box with safety rules found on slide 2")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Safety rules font size is ~20pt / 254000 EMU (0.15 points)
    try:
        if rules_box is not None:
            sizes_correct = 0
            total_runs = 0
            for para in rules_box.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        total_runs += 1
                        if run.font.size is not None:
                            # 20pt = 254000 EMU; allow small tolerance
                            if abs(run.font.size - 254000) <= 12700:  # +-1pt tolerance
                                sizes_correct += 1
                            else:
                                print(f"  SIZE MISMATCH: '{run.text[:30]}' size={run.font.size} EMU ({run.font.size/12700:.1f}pt)")

            if total_runs > 0 and sizes_correct == total_runs:
                print(f"PASS: Component 4 — All {total_runs} runs at ~20pt (0.15 pts)")
                total_score += 0.15
            elif total_runs > 0 and sizes_correct > 0:
                partial = round(0.15 * (sizes_correct / total_runs), 2)
                print(f"PARTIAL: Component 4 — {sizes_correct}/{total_runs} runs at correct size ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 4 — No runs at ~20pt font size")
        else:
            print(f"FAIL: Component 4 — No rules text box to check font size")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Triangle is positioned above the rules text box (0.15 points)
    try:
        if triangle is not None and rules_box is not None:
            tri_bottom = triangle.top + triangle.height
            rules_top = rules_box.top
            if triangle.top < rules_top:
                print(f"PASS: Component 5 — Triangle top ({triangle.top}) is above rules top ({rules_top}) (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 5 — Triangle top ({triangle.top}) is NOT above rules top ({rules_top})")
        else:
            missing = []
            if triangle is None:
                missing.append("triangle")
            if rules_box is None:
                missing.append("rules textbox")
            print(f"FAIL: Component 5 — Missing: {', '.join(missing)}")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
