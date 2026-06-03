"""
Initial Setup: Create a 5-slide Lab Safety presentation with slide 2 titled 'Safety First!' but empty.
Task ID: impress_teach_086
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_086'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Laboratory Safety Training"
    slide1.placeholders[1].text = "Department of Chemistry\nAcademic Year 2025-2026"

    # --- Slide 2: Safety First! (empty body - task target) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add only the title as a textbox
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Safety First!"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 3: Emergency Contacts ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    txTitle3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf3 = txTitle3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Emergency Contacts"
    p3.alignment = PP_ALIGN.LEFT
    r3 = p3.runs[0]
    r3.font.name = "Arial"
    r3.font.size = Pt(32)
    r3.font.bold = True

    txBody3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8), Inches(4.5))
    tf3b = txBody3.text_frame
    tf3b.word_wrap = True
    contacts = [
        "Campus Security: (555) 234-5678",
        "Lab Supervisor: Dr. Angela Reyes, ext. 4102",
        "Poison Control: 1-800-222-1222",
        "Facilities Emergency: (555) 234-9999",
        "Department Chair: Dr. Robert Kim, ext. 4001",
    ]
    for i, c in enumerate(contacts):
        if i == 0:
            p = tf3b.paragraphs[0]
        else:
            p = tf3b.add_paragraph()
        p.text = c
        p.space_after = Pt(8)
        r = p.runs[0]
        r.font.name = "Arial"
        r.font.size = Pt(18)

    # --- Slide 4: Required PPE ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txTitle4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf4 = txTitle4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Required Personal Protective Equipment"
    p4.alignment = PP_ALIGN.LEFT
    r4 = p4.runs[0]
    r4.font.name = "Arial"
    r4.font.size = Pt(28)
    r4.font.bold = True

    txBody4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8), Inches(4.5))
    tf4b = txBody4.text_frame
    tf4b.word_wrap = True
    ppe_items = [
        "Safety goggles (ANSI Z87.1 rated)",
        "Lab coat (knee-length, cotton or flame-resistant)",
        "Closed-toe shoes (no sandals or open heels)",
        "Chemical-resistant gloves (nitrile recommended)",
        "Face shield for splash hazards",
    ]
    for i, item in enumerate(ppe_items):
        if i == 0:
            p = tf4b.paragraphs[0]
        else:
            p = tf4b.add_paragraph()
        p.text = item
        p.space_after = Pt(8)
        r = p.runs[0]
        r.font.name = "Arial"
        r.font.size = Pt(18)

    # --- Slide 5: Lab Schedule ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    txTitle5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf5 = txTitle5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Weekly Lab Schedule"
    p5.alignment = PP_ALIGN.LEFT
    r5 = p5.runs[0]
    r5.font.name = "Arial"
    r5.font.size = Pt(32)
    r5.font.bold = True

    # Add a table for the schedule
    table_shape = slide5.shapes.add_table(4, 3, Inches(0.8), Inches(1.5), Inches(7.5), Inches(3))
    table = table_shape.table
    headers = ["Day", "Time", "Section"]
    data_rows = [
        ["Monday", "9:00 AM - 12:00 PM", "Section A - Organic Chemistry"],
        ["Wednesday", "1:00 PM - 4:00 PM", "Section B - Analytical Chemistry"],
        ["Friday", "10:00 AM - 1:00 PM", "Section C - Physical Chemistry"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(16)
    for r, row_data in enumerate(data_rows, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(14)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
