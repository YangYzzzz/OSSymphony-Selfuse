"""
Reward Script: Two grouped bar charts on slide 6 — Revenue by Region and Growth Rate by Region
Task ID: impress_exec_034
Domain: libreoffice_impress
Scoring:
  Component 1 (0.2): Two chart shapes exist on slide 6
  Component 2 (0.4): Left chart — title, position (~0.5in), type, categories, data values
  Component 3 (0.4): Right chart — title, position (~7in), type, categories, data values
"""

import os
from pptx import Presentation
from pptx.util import Inches

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_034'

# Expected data
EXPECTED_REVENUE_VALUES = [25.0, 18.0, 12.0, 7.5]
EXPECTED_GROWTH_VALUES = [12.0, 18.0, 35.0, 22.0]
EXPECTED_CATEGORIES = ['North America', 'Europe', 'APAC', 'LATAM']

# Position tolerance: left chart ~0.5in = 457200 EMU, right chart ~7in = 6400800 EMU
LEFT_CHART_X_TARGET = Inches(0.5)   # 457200
RIGHT_CHART_X_TARGET = Inches(7.0)  # 6400800
POSITION_TOLERANCE = Inches(0.5)    # generous tolerance for positioning


def get_chart_shapes(slide):
    """Return list of shapes that contain charts."""
    return [s for s in slide.shapes if hasattr(s, 'chart')]


def verify_chart(chart_shape, expected_title, expected_values, expected_cats, x_target):
    """Verify a single chart. Returns (sub_score out of 1.0, details_list)."""
    details = []
    sub = 0.0

    ch = chart_shape.chart

    # Check chart type is column/bar clustered (51 = COLUMN_CLUSTERED)
    if ch.chart_type == 51:
        sub += 0.15
        details.append("chart_type=COLUMN_CLUSTERED OK")
    else:
        details.append(f"chart_type={ch.chart_type}, expected COLUMN_CLUSTERED (51)")

    # Check title
    if ch.has_title:
        actual_title = ch.chart_title.text_frame.text.strip()
        if actual_title == expected_title:
            sub += 0.25
            details.append(f"title='{actual_title}' OK")
        else:
            details.append(f"title='{actual_title}', expected '{expected_title}'")
    else:
        details.append("no chart title")

    # Check position (x coordinate)
    actual_x = chart_shape.left
    if abs(actual_x - x_target) <= POSITION_TOLERANCE:
        sub += 0.15
        details.append(f"x={actual_x/914400:.2f}in OK (target {x_target/914400:.2f}in)")
    else:
        details.append(f"x={actual_x/914400:.2f}in, expected ~{x_target/914400:.2f}in")

    # Check categories
    try:
        plot = ch.plots[0]
        actual_cats = [str(c) for c in plot.categories]
        if actual_cats == expected_cats:
            sub += 0.20
            details.append(f"categories={actual_cats} OK")
        else:
            details.append(f"categories={actual_cats}, expected {expected_cats}")
    except Exception as e:
        details.append(f"categories error: {e}")

    # Check data values
    try:
        actual_values = list(ch.series[0].values)
        if actual_values == expected_values:
            sub += 0.25
            details.append(f"values={actual_values} OK")
        else:
            # Allow close float comparison
            if len(actual_values) == len(expected_values) and all(
                abs(a - e) < 0.01 for a, e in zip(actual_values, expected_values)
            ):
                sub += 0.25
                details.append(f"values={actual_values} OK (approx match)")
            else:
                details.append(f"values={actual_values}, expected {expected_values}")
    except Exception as e:
        details.append(f"values error: {e}")

    return sub, details


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide6 = prs.slides[5]
    chart_shapes = get_chart_shapes(slide6)

    # Component 1: Two chart shapes exist on slide 6 (0.2 points)
    try:
        num_charts = len(chart_shapes)
        if num_charts >= 2:
            print(f"PASS: Component 1 — {num_charts} charts on slide 6 (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 1 — {num_charts} charts on slide 6, expected >= 2")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    if len(chart_shapes) < 2:
        # Cannot proceed without at least 2 charts
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Sort charts by x position so left is first
    chart_shapes_sorted = sorted(chart_shapes, key=lambda s: s.left)
    left_chart = chart_shapes_sorted[0]
    right_chart = chart_shapes_sorted[1]

    # Component 2: Left chart — Revenue by Region (0.4 points)
    try:
        sub_score, details = verify_chart(
            left_chart, "Revenue by Region",
            EXPECTED_REVENUE_VALUES, EXPECTED_CATEGORIES,
            LEFT_CHART_X_TARGET
        )
        comp2_score = round(sub_score * 0.4, 4)
        if comp2_score > 0:
            total_score += comp2_score
        status = "PASS" if sub_score >= 0.99 else "PARTIAL"
        print(f"{status}: Component 2 — Left chart ({comp2_score:.2f}/0.40 pts)")
        for d in details:
            print(f"    {d}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Right chart — Growth Rate by Region (0.4 points)
    try:
        sub_score, details = verify_chart(
            right_chart, "Growth Rate by Region",
            EXPECTED_GROWTH_VALUES, EXPECTED_CATEGORIES,
            RIGHT_CHART_X_TARGET
        )
        comp3_score = round(sub_score * 0.4, 4)
        if comp3_score > 0:
            total_score += comp3_score
        status = "PASS" if sub_score >= 0.99 else "PARTIAL"
        print(f"{status}: Component 3 — Right chart ({comp3_score:.2f}/0.40 pts)")
        for d in details:
            print(f"    {d}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
