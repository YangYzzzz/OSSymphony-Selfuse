"""
Initial Setup: Create a 9-slide Regional Performance Review presentation.
Task ID: impress_exec_034
Domain: libreoffice_impress
Slide 6 has title 'Regional Performance' with no content (no charts).
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_034'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_body_slide(prs, layout_idx, title_text, body_lines):
    """Helper to add a slide with a title and bulleted body content."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title_text
    if body_lines and len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, line in enumerate(body_lines):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                p = tf.add_paragraph()
                p.text = line
    return slide


def create_initial():
    prs = Presentation()
    # Standard widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Regional Performance Review"
    slide1.placeholders[1].text = "Q4 2025 — Global Operations Division"

    # --- Slide 2: Market Overview ---
    add_title_body_slide(prs, 1, "Market Overview", [
        "Global market grew 8.2% year-over-year in Q4 2025",
        "North America maintained dominant position with 42% market share",
        "APAC region showed strongest acceleration at 35% growth",
        "European markets stabilized after regulatory adjustments",
        "Latin America emerging as key growth corridor",
    ])

    # --- Slide 3: Financial Summary ---
    add_title_body_slide(prs, 1, "Financial Summary", [
        "Total revenue: $62.5M across all regions",
        "Operating margin improved to 18.3% from 15.7%",
        "Cost optimization initiatives saved $4.2M in Q4",
        "Capital expenditure aligned with annual targets",
        "Cash reserves at $128M, up 12% from Q3",
    ])

    # --- Slide 4: Team Updates ---
    add_title_body_slide(prs, 1, "Team Updates", [
        "Headcount increased by 47 across engineering and sales",
        "Employee satisfaction score: 4.3/5.0 (up from 4.1)",
        "Completed leadership training program for 23 managers",
        "New regional director appointed for APAC operations",
        "Remote work policy updated for all global offices",
    ])

    # --- Slide 5: Strategic Priorities ---
    add_title_body_slide(prs, 1, "Strategic Priorities", [
        "Expand APAC presence with new Singapore hub",
        "Launch enterprise tier product by Q2 2026",
        "Achieve SOC 2 Type II certification",
        "Grow European channel partner network by 40%",
        "Reduce customer acquisition cost by 15%",
    ])

    # --- Slide 6: Regional Performance (TITLE ONLY, no charts) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Regional Performance"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 7: Action Items ---
    add_title_body_slide(prs, 1, "Action Items", [
        "Finalize APAC expansion budget by January 15",
        "Complete vendor selection for CRM migration",
        "Submit regulatory compliance documents for EU operations",
        "Schedule quarterly business reviews with regional leads",
        "Prepare investor update materials for board meeting",
    ])

    # --- Slide 8: Timeline ---
    add_title_body_slide(prs, 1, "Timeline & Milestones", [
        "Jan 2026: APAC hub operational planning complete",
        "Feb 2026: Enterprise product beta launch",
        "Mar 2026: SOC 2 audit begins",
        "Apr 2026: European partner summit",
        "Jun 2026: Mid-year performance review",
    ])

    # --- Slide 9: Thank You ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[0])
    slide9.shapes.title.text = "Thank You"
    slide9.placeholders[1].text = "Questions? Contact: global-ops@company.com"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
