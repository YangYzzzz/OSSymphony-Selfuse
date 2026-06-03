"""
Initial Setup: Financial report presentation with 8 slides.
Slide 5 titled 'Revenue Trends' has only title text.
Also creates /home/user/charts/revenue_graph.png image.
Task ID: impress_tm_090
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_090'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
CHART_DIR = f'{WORKDIR}/charts'
CHART_IMG = f'{CHART_DIR}/revenue_graph.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_chart_image():
    """Create a realistic-looking revenue bar chart image."""
    os.makedirs(CHART_DIR, exist_ok=True)

    width, height = 800, 500
    img = Image.new('RGB', (width, height), '#FFFFFF')
    draw = ImageDraw.Draw(img)

    # Background
    draw.rectangle([50, 30, 760, 440], outline='#333333', width=1)

    # Title
    draw.text((250, 8), "Annual Revenue (2020-2025)", fill='#333333')

    # Y-axis labels and grid lines
    revenues = [120, 145, 160, 185, 210, 248]
    max_rev = 280
    bar_colors = ['#4472C4', '#4472C4', '#4472C4', '#4472C4', '#4472C4', '#ED7D31']

    for i in range(0, 7):
        y_val = i * 40
        y_pos = 440 - int((y_val / max_rev) * 410)
        draw.line([(50, y_pos), (760, y_pos)], fill='#DDDDDD', width=1)
        draw.text((10, y_pos - 6), f"${y_val}M", fill='#666666')

    # Bars
    bar_width = 80
    gap = 30
    start_x = 100
    years = ['2020', '2021', '2022', '2023', '2024', '2025']

    for i, (rev, year) in enumerate(zip(revenues, years)):
        x1 = start_x + i * (bar_width + gap)
        x2 = x1 + bar_width
        bar_height = int((rev / max_rev) * 410)
        y1 = 440 - bar_height
        y2 = 440

        draw.rectangle([x1, y1, x2, y2], fill=bar_colors[i])
        draw.text((x1 + 15, 445), year, fill='#333333')
        draw.text((x1 + 10, y1 - 16), f"${rev}M", fill='#333333')

    img.save(CHART_IMG)
    print(f'Chart image created: {CHART_IMG}')


def create_initial():
    prs = Presentation()

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Meridian Technologies Inc."
    slide1.placeholders[1].text = "Annual Financial Report 2025"

    # Slide 2: Agenda
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Executive Summary"
    for item in ["Financial Highlights", "Market Analysis",
                 "Revenue Trends", "Regional Breakdown",
                 "Strategic Outlook", "Q&A"]:
        p = body2.add_paragraph()
        p.text = item
        p.level = 0

    # Slide 3: Executive Summary
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Executive Summary"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Meridian Technologies delivered exceptional results in FY2025, with revenue growth of 18% year-over-year driven by strong enterprise demand."
    p3 = body3.add_paragraph()
    p3.text = ""
    p3b = body3.add_paragraph()
    p3b.text = "Key achievements include expansion into 12 new markets, acquisition of DataFlow Analytics, and the launch of our AI-powered automation platform."

    # Slide 4: Financial Highlights
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    txBox4_title = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf4_title = txBox4_title.text_frame
    p4t = tf4_title.paragraphs[0]
    p4t.text = "Financial Highlights"
    run4t = p4t.runs[0]
    run4t.font.size = Pt(28)
    run4t.font.bold = True
    run4t.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Add a simple table for financial data
    table_shape = slide4.shapes.add_table(5, 3, Inches(1), Inches(1.5), Inches(7), Inches(3))
    table = table_shape.table
    headers = ["Metric", "FY2024", "FY2025"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
    data4 = [
        ["Total Revenue", "$210M", "$248M"],
        ["Gross Margin", "62.3%", "64.8%"],
        ["Operating Income", "$48.5M", "$62.1M"],
        ["Net Income", "$35.2M", "$47.8M"],
    ]
    for r, row in enumerate(data4, 1):
        for c, val in enumerate(row):
            table.cell(r, c).text = val

    # Slide 5: Revenue Trends (ONLY title, no image, no caption)
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    txBox5_title = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf5_title = txBox5_title.text_frame
    p5t = tf5_title.paragraphs[0]
    p5t.text = "Revenue Trends"
    run5t = p5t.runs[0]
    run5t.font.size = Pt(28)
    run5t.font.bold = True
    run5t.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Slide 6: Regional Breakdown
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Regional Revenue Breakdown"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "North America: $142M (57%)"
    for item in ["Europe: $58M (23%)", "Asia-Pacific: $35M (14%)", "Rest of World: $13M (6%)"]:
        p6 = body6.add_paragraph()
        p6.text = item

    # Slide 7: Strategic Outlook
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Strategic Outlook for 2026"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "Target 22% revenue growth through expanded product portfolio"
    for item in ["Planned investment of $30M in R&D for AI capabilities",
                 "Geographic expansion into Latin America and Middle East",
                 "Strategic partnerships with three Fortune 100 companies"]:
        p7 = body7.add_paragraph()
        p7.text = item

    # Slide 8: Q&A
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    txBox8 = slide8.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(2))
    tf8 = txBox8.text_frame
    p8 = tf8.paragraphs[0]
    p8.text = "Questions & Discussion"
    p8.alignment = PP_ALIGN.CENTER
    run8 = p8.runs[0]
    run8.font.size = Pt(36)
    run8.font.bold = True
    run8.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    p8b = tf8.add_paragraph()
    p8b.text = "Thank you for your attention"
    p8b.alignment = PP_ALIGN.CENTER
    run8b = p8b.runs[0]
    run8b.font.size = Pt(18)
    run8b.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_chart_image()
create_initial()
