"""
Reward Script: Insert revenue graph image on slide 5 with caption text box
Task ID: impress_tm_090
Domain: libreoffice_impress
Scoring:
  Component 1 (0.35): Image present on slide 5 matching revenue_graph.png
  Component 2 (0.25): Caption text box with correct text on slide 5
  Component 3 (0.20): Caption font is italic and 10pt
  Component 4 (0.20): Caption font color is gray (#808080)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_090'
IMAGE_PATH = '/home/user/charts/revenue_graph.png'
EXPECTED_CAPTION = 'Figure 1: Annual Revenue (2020-2025)'


def persist_app_state():
    """Send Ctrl+S to save any unsaved LibreOffice edits."""
    try:
        os.environ["DISPLAY"] = ":0"
        import pyautogui
        import time
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide5 = prs.slides[4]  # 0-indexed

    # Load reference image blob for comparison
    ref_blob = None
    try:
        with open(IMAGE_PATH, 'rb') as f:
            ref_blob = f.read()
    except Exception as e:
        print(f"WARN: Could not load reference image {IMAGE_PATH}: {e}")

    # Component 1: Image present on slide 5 matching revenue_graph.png (0.35 points)
    try:
        image_found = False
        for shape in slide5.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if ref_blob is not None:
                    if shape.image.blob == ref_blob:
                        image_found = True
                        print(f"PASS: Component 1 — revenue_graph.png image found on slide 5 (blob match, {len(shape.image.blob)} bytes) (0.35 pts)")
                        total_score += 0.35
                        break
                else:
                    # Fallback: any image on slide 5 counts if we can't load reference
                    image_found = True
                    print(f"PASS: Component 1 — Image found on slide 5 (no ref to compare, blob_size={len(shape.image.blob)}) (0.35 pts)")
                    total_score += 0.35
                    break
        if not image_found:
            # Check if there's any picture at all (wrong image)
            pics = [s for s in slide5.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
            if pics:
                print(f"FAIL: Component 1 — Image found on slide 5 but blob does not match revenue_graph.png")
            else:
                print(f"FAIL: Component 1 — No image found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Find caption text box on slide 5
    caption_shape = None
    caption_text = None
    for shape in slide5.shapes:
        if shape.has_text_frame and shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            full_text = shape.text_frame.text.strip()
            # Skip the "Revenue Trends" title text box
            if full_text and 'Revenue Trends' not in full_text:
                caption_shape = shape
                caption_text = full_text
                break

    # Component 2: Caption text box with correct text (0.25 points)
    try:
        if caption_shape is not None and caption_text == EXPECTED_CAPTION:
            print(f"PASS: Component 2 — Caption text matches exactly: '{caption_text}' (0.25 pts)")
            total_score += 0.25
        elif caption_shape is not None and EXPECTED_CAPTION.lower() in caption_text.lower():
            print(f"PARTIAL: Component 2 — Caption text partially matches: '{caption_text}' (0.15 pts)")
            total_score += 0.15
        elif caption_shape is not None:
            print(f"FAIL: Component 2 — Caption text mismatch. Expected: '{EXPECTED_CAPTION}', Found: '{caption_text}'")
        else:
            print(f"FAIL: Component 2 — No caption text box found on slide 5 (only title text box present)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Caption font is italic and 10pt (0.20 points)
    try:
        if caption_shape is not None:
            runs = []
            for para in caption_shape.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        runs.append(run)

            if runs:
                run = runs[0]
                is_italic = run.font.italic is True
                # 10pt = 127000 EMU
                is_10pt = run.font.size is not None and run.font.size == 127000

                if is_italic and is_10pt:
                    print(f"PASS: Component 3 — Caption font is italic={run.font.italic}, size={run.font.size} EMU (10pt) (0.20 pts)")
                    total_score += 0.20
                elif is_italic:
                    print(f"PARTIAL: Component 3 — Italic correct but size={run.font.size} EMU (expected 127000/10pt) (0.10 pts)")
                    total_score += 0.10
                elif is_10pt:
                    print(f"PARTIAL: Component 3 — Size correct (10pt) but italic={run.font.italic} (expected True) (0.10 pts)")
                    total_score += 0.10
                else:
                    print(f"FAIL: Component 3 — italic={run.font.italic} (expected True), size={run.font.size} (expected 127000)")
            else:
                print(f"FAIL: Component 3 — No runs found in caption text box")
        else:
            print(f"FAIL: Component 3 — No caption text box to check font")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Caption font color is gray #808080 (0.20 points)
    try:
        if caption_shape is not None:
            runs = []
            for para in caption_shape.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        runs.append(run)

            if runs:
                run = runs[0]
                color_match = False
                try:
                    if run.font.color.type is not None:
                        rgb = str(run.font.color.rgb).upper()
                        if rgb == '808080':
                            color_match = True
                            print(f"PASS: Component 4 — Caption font color is #{rgb} (gray) (0.20 pts)")
                            total_score += 0.20
                        else:
                            print(f"FAIL: Component 4 — Caption font color is #{rgb}, expected #808080")
                    else:
                        print(f"FAIL: Component 4 — Caption font color type is None (no explicit color set)")
                except AttributeError:
                    print(f"FAIL: Component 4 — Could not read caption font color (theme-based or not set)")
            else:
                print(f"FAIL: Component 4 — No runs found in caption text box")
        else:
            print(f"FAIL: Component 4 — No caption text box to check color")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
