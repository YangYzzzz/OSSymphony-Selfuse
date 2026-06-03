"""
Initial Setup: Add tab stops to text box on slide 4
Task ID: impress_tct_094
Domain: libreoffice_impress

Creates a 5-slide presentation 'Tabbed_Layout.pptx'.
Slide 4 has a text box with tab characters but NO defined tab stops
(text uses default tab spacing).
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_094'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title Slide ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q3 2025 Regional Sales Report"
    slide1.placeholders[1].text = "Prepared by Analytics Division\nSeptember 2025"

    # ── Slide 2: Executive Summary ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Total revenue reached $4.87M across all regions in Q3."
    p = body2.add_paragraph()
    p.text = "West Coast led with 38% of total revenue, followed by Midwest at 27%."
    p.level = 0
    p = body2.add_paragraph()
    p.text = "Customer acquisition cost decreased by 12% compared to Q2."
    p.level = 0
    p = body2.add_paragraph()
    p.text = "Net promoter score improved to 72, up from 64 in Q2."
    p.level = 0

    # ── Slide 3: Revenue Breakdown ──
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf3 = title3.text_frame
    p = tf3.paragraphs[0]
    p.text = "Revenue Breakdown by Region"
    p.runs[0].font.size = Pt(28)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    table_shape = slide3.shapes.add_table(
        6, 4, Inches(1), Inches(1.5), Inches(10), Inches(3.5)
    )
    table = table_shape.table
    headers = ["Region", "Revenue ($K)", "Growth (%)", "Target Hit"]
    data = [
        ["West Coast", "1,851", "+14.2", "Yes"],
        ["Midwest", "1,315", "+8.7", "Yes"],
        ["Southeast", "892", "+3.1", "No"],
        ["Northeast", "618", "-1.4", "No"],
        ["Mountain", "194", "+22.6", "Yes"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row in enumerate(data, 1):
        for c, val in enumerate(row):
            table.cell(r, c).text = val

    # ── Slide 4: Tabbed Layout Data ──
    # This slide has a text box with TAB characters but NO custom tab stops.
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.8))
    tf4t = title4.text_frame
    p = tf4t.paragraphs[0]
    p.text = "Employee Compensation Summary"
    p.runs[0].font.size = Pt(28)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Main text box with tab-separated content (NO tab stops defined)
    txBox = slide4.shapes.add_textbox(
        Inches(0.8), Inches(1.4), Inches(10), Inches(5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    # Header line with tabs
    p_hdr = tf.paragraphs[0]
    run_hdr = p_hdr.runs[0] if p_hdr.runs else p_hdr.add_run()
    run_hdr.text = "Name\tDepartment\tAnnual Salary"
    run_hdr.font.bold = True
    run_hdr.font.size = Pt(16)
    run_hdr.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Data lines with tabs
    employees = [
        ("Sarah Chen", "Engineering", "$128,500"),
        ("Marcus Johnson", "Marketing", "$97,200"),
        ("Priya Patel", "Data Science", "$115,800"),
        ("David Kim", "Product", "$108,400"),
        ("Elena Rodriguez", "Design", "$92,600"),
        ("James Thompson", "Finance", "$104,300"),
        ("Aisha Williams", "Operations", "$89,750"),
        ("Robert Garcia", "Engineering", "$132,000"),
    ]
    for name, dept, salary in employees:
        p_data = tf.add_paragraph()
        run_data = p_data.add_run()
        run_data.text = f"{name}\t{dept}\t{salary}"
        run_data.font.size = Pt(14)
        run_data.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # ── Slide 5: Closing / Next Steps ──
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Next Steps"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Finalize compensation adjustments by October 15, 2025."
    p = body5.add_paragraph()
    p.text = "Submit regional budget proposals for Q4 planning."
    p = body5.add_paragraph()
    p.text = "Schedule leadership review for November 3, 2025."

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
