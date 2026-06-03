"""
Reward Script: Verify 'Photo Feature' custom layout creation and application
Task ID: impress_rp_014
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): 'Photo Feature' layout exists in the presentation
  Component 2 (0.25): Layout has correct placeholder structure (picture left 40%, title top-right, body bottom-right)
  Component 3 (0.30): Slides 3, 5, 7 use the 'Photo Feature' layout
  Component 4 (0.20): Other slides (1,2,4,6,8,9,10) do NOT use 'Photo Feature' layout
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_014'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slide_width = prs.slide_width  # 9144000 EMU = 10 inches
    slide_height = prs.slide_height  # 6858000 EMU = 7.5 inches

    # Component 1: 'Photo Feature' layout exists (0.25 points)
    photo_feature_layout = None
    try:
        for layout in prs.slide_layouts:
            if layout.name == 'Photo Feature':
                photo_feature_layout = layout
                break

        if photo_feature_layout is not None:
            print(f"PASS: Component 1 — 'Photo Feature' layout exists (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — 'Photo Feature' layout not found. Available layouts: {[l.name for l in prs.slide_layouts]}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Layout has correct placeholder structure (0.25 points)
    # Expected: picture placeholder on left (40% width, full height),
    #           title top-right (from 40% to 100% width, top half),
    #           body bottom-right (from 40% to 100% width, bottom half)
    try:
        if photo_feature_layout is not None:
            placeholders = list(photo_feature_layout.placeholders)
            # Filter out date/footer/slide number placeholders (idx 10,11,12 for standard ones)
            # But note: in the golden file, the picture placeholder uses idx=10
            # We need to check by type, not by index

            has_picture = False
            has_title = False
            has_body = False
            picture_ok = False
            title_ok = False
            body_ok = False

            width_40pct = int(slide_width * 0.4)  # 3657600
            height_50pct = int(slide_height * 0.5)  # 3429000
            tolerance = 0.05  # 5% tolerance for position checks

            for ph in placeholders:
                ph_type = ph.placeholder_format.type
                # Check for PICTURE type (18)
                if ph_type == 18:  # PICTURE
                    has_picture = True
                    # Should be on left: left ~0, width ~40% of slide, height ~full
                    left_ok = ph.left <= int(slide_width * tolerance)
                    width_ok = abs(ph.width - width_40pct) <= int(slide_width * tolerance)
                    height_ok = abs(ph.height - slide_height) <= int(slide_height * tolerance)
                    picture_ok = left_ok and width_ok and height_ok
                    if picture_ok:
                        print(f"  Picture placeholder: left={ph.left}, width={ph.width}, height={ph.height} — OK")
                    else:
                        print(f"  Picture placeholder: left={ph.left}, width={ph.width}, height={ph.height} — position/size mismatch")
                        print(f"    Expected: left~0, width~{width_40pct}, height~{slide_height}")

                # Check for TITLE type (1 or 3)
                elif ph_type in (1, 3):  # TITLE or CENTER_TITLE
                    has_title = True
                    # Should be top-right: left ~40%, top ~0, width ~60%, height ~50%
                    left_ok = abs(ph.left - width_40pct) <= int(slide_width * tolerance)
                    top_ok = ph.top <= int(slide_height * tolerance)
                    width_ok = abs(ph.width - (slide_width - width_40pct)) <= int(slide_width * tolerance)
                    height_ok = abs(ph.height - height_50pct) <= int(slide_height * tolerance)
                    title_ok = left_ok and top_ok and width_ok and height_ok
                    if title_ok:
                        print(f"  Title placeholder: left={ph.left}, top={ph.top}, width={ph.width}, height={ph.height} — OK")
                    else:
                        print(f"  Title placeholder: left={ph.left}, top={ph.top}, width={ph.width}, height={ph.height} — position/size mismatch")

                # Check for BODY type (2)
                elif ph_type == 2:  # BODY
                    has_body = True
                    # Should be bottom-right: left ~40%, top ~50%, width ~60%, height ~50%
                    left_ok = abs(ph.left - width_40pct) <= int(slide_width * tolerance)
                    top_ok = abs(ph.top - height_50pct) <= int(slide_height * tolerance)
                    width_ok = abs(ph.width - (slide_width - width_40pct)) <= int(slide_width * tolerance)
                    height_ok = abs(ph.height - height_50pct) <= int(slide_height * tolerance)
                    body_ok = left_ok and top_ok and width_ok and height_ok
                    if body_ok:
                        print(f"  Body placeholder: left={ph.left}, top={ph.top}, width={ph.width}, height={ph.height} — OK")
                    else:
                        print(f"  Body placeholder: left={ph.left}, top={ph.top}, width={ph.width}, height={ph.height} — position/size mismatch")

            sub_score = 0.0
            if has_picture and picture_ok:
                sub_score += 1/3
            if has_title and title_ok:
                sub_score += 1/3
            if has_body and body_ok:
                sub_score += 1/3

            component_2_score = round(0.25 * sub_score, 4)
            if component_2_score > 0:
                print(f"PASS: Component 2 — Layout placeholder structure verified ({component_2_score:.4f} pts)")
                total_score += component_2_score
            else:
                print(f"FAIL: Component 2 — Layout placeholder structure incorrect. picture={has_picture}/{picture_ok}, title={has_title}/{title_ok}, body={has_body}/{body_ok}")
        else:
            print(f"FAIL: Component 2 — Cannot check layout structure (no 'Photo Feature' layout found)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slides 3, 5, 7 use 'Photo Feature' layout (0.30 points)
    slides_using_pf = 0
    try:
        target_slides = [2, 4, 6]  # 0-indexed: slides 3, 5, 7

        for idx in target_slides:
            if idx < len(prs.slides):
                slide = prs.slides[idx]
                layout_name = slide.slide_layout.name
                if layout_name == 'Photo Feature':
                    slides_using_pf += 1
                    print(f"  Slide {idx+1} layout: '{layout_name}' — OK")
                else:
                    print(f"  Slide {idx+1} layout: '{layout_name}' — expected 'Photo Feature'")
            else:
                print(f"  Slide {idx+1}: does not exist")

        if slides_using_pf == 3:
            print(f"PASS: Component 3 — All 3 target slides use 'Photo Feature' (0.30 pts)")
            total_score += 0.30
        elif slides_using_pf > 0:
            partial = round(0.30 * (slides_using_pf / 3), 4)
            print(f"PARTIAL: Component 3 — {slides_using_pf}/3 target slides use 'Photo Feature' ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No target slides use 'Photo Feature'")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Other slides do NOT use 'Photo Feature' AND target slides DO (0.20 points)
    # Gate: only scores if 'Photo Feature' layout exists AND at least one target slide uses it
    # This ensures we are scoring a task-introduced change (selective application)
    try:
        if photo_feature_layout is not None and slides_using_pf > 0:
            other_slides = [0, 1, 3, 5, 7, 8, 9]  # 0-indexed: slides 1,2,4,6,8,9,10
            others_correct = 0

            for idx in other_slides:
                if idx < len(prs.slides):
                    slide = prs.slides[idx]
                    layout_name = slide.slide_layout.name
                    if layout_name != 'Photo Feature':
                        others_correct += 1
                    else:
                        print(f"  Slide {idx+1} layout: '{layout_name}' — should NOT be 'Photo Feature'")

            if others_correct == len(other_slides):
                print(f"PASS: Component 4 — Other slides do not use 'Photo Feature' (0.20 pts)")
                total_score += 0.20
            else:
                wrong_count = len(other_slides) - others_correct
                print(f"FAIL: Component 4 — {wrong_count} non-target slides incorrectly use 'Photo Feature'")
        else:
            print(f"FAIL: Component 4 — Precondition not met (no 'Photo Feature' layout or no target slides use it)")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
