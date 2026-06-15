"""
Initial Setup: Create a 10-slide portfolio presentation with blank layouts.
Task ID: impress_rp_014
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_014'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=None):
    """Helper to add a text box with formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()

    # Standard slide dimensions (10 x 7.5 inches)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Use blank layout (index 6 in default template)
    blank_layout = prs.slide_layouts[6]  # Blank layout

    # Slide 1: Title slide (using blank layout with text boxes)
    slide1 = prs.slides.add_slide(blank_layout)
    add_textbox(slide1, Inches(1), Inches(1.5), Inches(8), Inches(1.5),
                "Creative Portfolio Showcase", font_size=36, bold=True,
                color=RGBColor(0x2E, 0x4A, 0x62), alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(2), Inches(3.5), Inches(6), Inches(1),
                "Elena Rodriguez | Visual Designer & Photographer", font_size=18,
                color=RGBColor(0x5B, 0x6E, 0x7E), alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(2.5), Inches(5), Inches(5), Inches(0.8),
                "Spring 2025 Collection", font_size=14,
                color=RGBColor(0x8E, 0x99, 0xA4), alignment=PP_ALIGN.CENTER)

    # Slide 2: About Me
    slide2 = prs.slides.add_slide(blank_layout)
    add_textbox(slide2, Inches(0.5), Inches(0.5), Inches(9), Inches(1),
                "About Me", font_size=28, bold=True,
                color=RGBColor(0x2E, 0x4A, 0x62))
    add_textbox(slide2, Inches(0.5), Inches(1.8), Inches(9), Inches(4),
                "With over 12 years of experience in visual design and photography, "
                "I bring a unique perspective to every project. My work spans commercial "
                "photography, brand identity design, and editorial layouts for publications "
                "including Vogue Italia, Architectural Digest, and National Geographic Traveler.",
                font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # Slide 3: Urban Architecture Series
    slide3 = prs.slides.add_slide(blank_layout)
    add_textbox(slide3, Inches(0.5), Inches(0.5), Inches(9), Inches(1),
                "Urban Architecture Series", font_size=28, bold=True,
                color=RGBColor(0x2E, 0x4A, 0x62))
    add_textbox(slide3, Inches(0.5), Inches(1.8), Inches(9), Inches(4),
                "Captured across 15 cities worldwide, this series explores the intersection "
                "of modern architecture and urban life. Featured locations include the Zaha Hadid "
                "Centre in Vienna, Marina Bay Sands in Singapore, and the CCTV Headquarters in Beijing.",
                font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # Slide 4: Client Work - Aether Brands
    slide4 = prs.slides.add_slide(blank_layout)
    add_textbox(slide4, Inches(0.5), Inches(0.5), Inches(9), Inches(1),
                "Client Work: Aether Brands", font_size=28, bold=True,
                color=RGBColor(0x2E, 0x4A, 0x62))
    add_textbox(slide4, Inches(0.5), Inches(1.8), Inches(9), Inches(2),
                "Complete brand identity redesign for Aether Brands, a sustainable fashion "
                "company based in Copenhagen. Deliverables included logo suite, typography system, "
                "packaging design, and social media templates. The project increased brand recognition "
                "by 47% within six months of launch.",
                font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # Slide 5: Nature & Wildlife Photography
    slide5 = prs.slides.add_slide(blank_layout)
    add_textbox(slide5, Inches(0.5), Inches(0.5), Inches(9), Inches(1),
                "Nature & Wildlife Photography", font_size=28, bold=True,
                color=RGBColor(0x2E, 0x4A, 0x62))
    add_textbox(slide5, Inches(0.5), Inches(1.8), Inches(9), Inches(4),
                "Award-winning wildlife photography from expeditions to Patagonia, "
                "the Serengeti, and the Norwegian fjords. This collection was featured in "
                "the 2024 World Press Photo exhibition and received the Environmental "
                "Photographer of the Year honorable mention.",
                font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # Slide 6: Editorial Work - Meridian Magazine
    slide6 = prs.slides.add_slide(blank_layout)
    add_textbox(slide6, Inches(0.5), Inches(0.5), Inches(9), Inches(1),
                "Editorial: Meridian Magazine", font_size=28, bold=True,
                color=RGBColor(0x2E, 0x4A, 0x62))
    add_textbox(slide6, Inches(0.5), Inches(1.8), Inches(9), Inches(2),
                "Monthly cover photography and feature spreads for Meridian Magazine "
                "since March 2023. Notable features include the 'Faces of Innovation' "
                "series profiling tech entrepreneurs across Southeast Asia, and the "
                "'Hidden Kitchens' culinary travel feature.",
                font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # Slide 7: Product Photography Portfolio
    slide7 = prs.slides.add_slide(blank_layout)
    add_textbox(slide7, Inches(0.5), Inches(0.5), Inches(9), Inches(1),
                "Product Photography Portfolio", font_size=28, bold=True,
                color=RGBColor(0x2E, 0x4A, 0x62))
    add_textbox(slide7, Inches(0.5), Inches(1.8), Inches(9), Inches(4),
                "Specialized product photography for luxury brands including Cartier, "
                "Bang & Olufsen, and Diptyque. Studio and lifestyle shots optimized for "
                "e-commerce platforms, print catalogs, and social media campaigns.",
                font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # Slide 8: Awards & Recognition
    slide8 = prs.slides.add_slide(blank_layout)
    add_textbox(slide8, Inches(0.5), Inches(0.5), Inches(9), Inches(1),
                "Awards & Recognition", font_size=28, bold=True,
                color=RGBColor(0x2E, 0x4A, 0x62))
    add_textbox(slide8, Inches(0.5), Inches(1.8), Inches(9), Inches(4),
                "- IPA International Photography Awards, Gold (2024)\n"
                "- Communication Arts Photography Annual (2023, 2024)\n"
                "- Red Dot Design Award, Brand Identity (2023)\n"
                "- PDN Photo Annual, Advertising Category (2022)\n"
                "- Graphis Design Annual, Platinum (2022)",
                font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # Slide 9: Exhibitions & Speaking
    slide9 = prs.slides.add_slide(blank_layout)
    add_textbox(slide9, Inches(0.5), Inches(0.5), Inches(9), Inches(1),
                "Exhibitions & Speaking", font_size=28, bold=True,
                color=RGBColor(0x2E, 0x4A, 0x62))
    add_textbox(slide9, Inches(0.5), Inches(1.8), Inches(9), Inches(3),
                "Solo exhibitions at Galerie Perrotin (Paris, 2024), The Photographers' "
                "Gallery (London, 2023), and ICP Museum (New York, 2023). Keynote speaker "
                "at Adobe MAX 2024 and Photo London 2024. Guest lecturer at Parsons School "
                "of Design and Royal College of Art.",
                font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # Slide 10: Contact
    slide10 = prs.slides.add_slide(blank_layout)
    add_textbox(slide10, Inches(1), Inches(1.5), Inches(8), Inches(1.2),
                "Let's Work Together", font_size=32, bold=True,
                color=RGBColor(0x2E, 0x4A, 0x62), alignment=PP_ALIGN.CENTER)
    add_textbox(slide10, Inches(2), Inches(3.2), Inches(6), Inches(3),
                "Elena Rodriguez\n"
                "elena@rodriguezcreative.com\n"
                "+1 (415) 892-3047\n"
                "www.rodriguezcreative.com\n"
                "@elena.rodriguez.photo",
                font_size=18, color=RGBColor(0x5B, 0x6E, 0x7E),
                alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Number of slides: {len(prs.slides)}')
    print(f'Number of slide layouts: {len(prs.slide_layouts)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
