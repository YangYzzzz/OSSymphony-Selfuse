"""
Reward Script: Reduce presentation file size by compressing images and removing unused masters
Task ID: impress_fix_026
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.40): File size under 40MB
  - Component 2 (0.35): Images compressed (avg blob < 500KB per image)
  - Component 3 (0.25): Unused master slides removed (count reduced to <= 2)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_026'


def persist_app_state(domain: str):
    """Save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            import time
            pyautogui.hotkey("ctrl", "s")
            time.sleep(1.0)
            print("PERSIST: ctrl+s sent for libreoffice_impress")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: file must exist
    if not os.path.exists(file_path):
        print(f"CRITICAL: File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file must be loadable
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition gate: slide count and content integrity
    # The task says "slide content and layout remain intact" — 30 slides expected
    if len(prs.slides) != 30:
        print(f"FAIL: Precondition — expected 30 slides, found {len(prs.slides)}. Content corrupted.")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: File size under 40MB (0.40 points)
    # Initial file is ~363MB, golden should be under 40MB per task context
    try:
        file_size = os.path.getsize(file_path)
        file_size_mb = file_size / (1024 * 1024)
        if file_size_mb < 40.0:
            print(f"PASS: Component 1 — File size {file_size_mb:.2f} MB is under 40MB target (0.40 pts)")
            total_score += 0.40
        elif file_size_mb < 120.0:
            # Partial credit: significant reduction but not fully meeting target
            partial = 0.20
            print(f"PARTIAL: Component 1 — File size {file_size_mb:.2f} MB reduced but above 40MB target ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — File size {file_size_mb:.2f} MB is not significantly reduced from original ~363MB")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Images compressed — average blob size < 500KB (0.35 points)
    # Initial images are ~12MB each; golden should be much smaller
    try:
        total_blob_size = 0
        image_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'image'):
                    try:
                        total_blob_size += len(shape.image.blob)
                        image_count += 1
                    except Exception:
                        pass

        if image_count == 0:
            print("FAIL: Component 2 — No images found; images should still be present after compression")
        else:
            avg_blob_kb = (total_blob_size / image_count) / 1024
            if avg_blob_kb < 500:
                print(f"PASS: Component 2 — Average image blob {avg_blob_kb:.1f} KB (< 500KB threshold), {image_count} images (0.35 pts)")
                total_score += 0.35
            elif avg_blob_kb < 2000:
                partial = 0.15
                print(f"PARTIAL: Component 2 — Average image blob {avg_blob_kb:.1f} KB partially compressed ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 2 — Average image blob {avg_blob_kb:.1f} KB, images not sufficiently compressed")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Unused master slides removed — count should be <= 2 (0.25 points)
    # Initial has 5 masters but only 2 are used; golden should have removed the 3 unused ones
    try:
        master_count = len(prs.slide_masters)
        if master_count <= 2:
            print(f"PASS: Component 3 — {master_count} master slide(s), unused masters removed (0.25 pts)")
            total_score += 0.25
        elif master_count < 5:
            partial = 0.10
            print(f"PARTIAL: Component 3 — {master_count} masters (reduced from 5 but not fully cleaned) ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — {master_count} master slides remain, expected <= 2 (unused masters not removed)")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
