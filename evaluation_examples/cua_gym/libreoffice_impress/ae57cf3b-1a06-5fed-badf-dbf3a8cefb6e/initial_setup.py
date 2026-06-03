"""
Initial Setup: Create a large presentation with high-resolution images and unused master slides
Task ID: impress_fix_026
Domain: libreoffice_impress

Creates Event_Recap.pptx with:
- 30 slides containing large images (high-res photos generated via Pillow)
- 5 slide masters defined (only 2 used by slides)
- Realistic event recap content
"""

import os
import shlex
import subprocess
import time
import io
import copy
import zipfile
import shutil
from lxml import etree

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_026'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def generate_large_image(width, height, base_color, text, quality='high'):
    """Generate a large, photo-like image with noise to ensure large PNG file size."""
    from PIL import Image, ImageDraw, ImageFont
    import random
    import numpy as np

    # Create base image with gradient
    r0, g0, b0 = base_color
    arr = np.zeros((height, width, 3), dtype=np.uint8)

    # Base gradient
    for y in range(height):
        ratio = y / height
        arr[y, :, 0] = int(r0 * (1 - ratio * 0.3))
        arr[y, :, 1] = int(g0 * (1 - ratio * 0.2))
        arr[y, :, 2] = int(b0 * (1 - ratio * 0.4))

    # Add photographic noise (makes PNG large since noise is incompressible)
    rng = np.random.RandomState(hash(text) % (2**31))
    noise = rng.randint(-25, 26, size=(height, width, 3), dtype=np.int16)
    arr = np.clip(arr.astype(np.int16) + noise, 0, 255).astype(np.uint8)

    img = Image.fromarray(arr)
    draw = ImageDraw.Draw(img)

    # Add random detail rectangles
    random.seed(hash(text) % 10000)
    for _ in range(60):
        x1 = random.randint(0, width - 100)
        y1 = random.randint(0, height - 100)
        x2 = x1 + random.randint(30, 200)
        y2 = y1 + random.randint(30, 200)
        alpha_r = random.randint(max(0, r0 - 70), min(255, r0 + 70))
        alpha_g = random.randint(max(0, g0 - 70), min(255, g0 + 70))
        alpha_b = random.randint(max(0, b0 - 70), min(255, b0 + 70))
        draw.rectangle([x1, y1, x2, y2], fill=(alpha_r, alpha_g, alpha_b))

    # Add text overlay with shadow
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 64)
    except (IOError, OSError):
        font = ImageFont.load_default()
    bbox = draw.textbbox((0, 0), text, font=font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    draw.text(((width - tw) // 2 + 2, (height - th) // 2 + 2), text, fill=(0, 0, 0), font=font)
    draw.text(((width - tw) // 2, (height - th) // 2), text, fill=(255, 255, 255), font=font)

    buf = io.BytesIO()
    # Save as PNG (lossless = large with noise)
    img.save(buf, format='PNG')
    buf.seek(0)
    return buf


def create_initial():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Event recap content for 30 slides
    slide_content = [
        ("Annual Tech Summit 2025", "Event Recap & Highlights", (41, 65, 122)),
        ("Event Overview", "March 15-17, 2025 | San Francisco Convention Center", (52, 73, 94)),
        ("Day 1: Opening Keynote", "CEO Sarah Mitchell's Vision for the Future", (70, 130, 180)),
        ("Registration & Welcome", "Over 2,500 attendees from 40 countries", (44, 62, 80)),
        ("Keynote Highlights", "AI-First Strategy Announcement", (30, 80, 120)),
        ("Panel: Future of Cloud Computing", "Industry Leaders Discuss Trends", (60, 90, 130)),
        ("Workshop: Machine Learning Basics", "Hands-on session with 200 participants", (80, 60, 90)),
        ("Networking Lunch - Day 1", "Grand Ballroom - International Cuisine", (120, 80, 40)),
        ("Product Demo: AutoML Platform", "Live demo with real customer data", (40, 100, 80)),
        ("Day 1 Evening Reception", "Rooftop Garden Party at The Mark Hotel", (80, 50, 100)),
        ("Day 2: Technical Deep Dives", "Advanced Sessions and Breakout Rooms", (50, 100, 120)),
        ("Session: Kubernetes at Scale", "Managing 10,000+ node clusters", (60, 70, 100)),
        ("Workshop: Data Pipeline Design", "Building resilient streaming architectures", (70, 90, 60)),
        ("Sponsor Showcase", "30+ technology partners demonstrated solutions", (90, 60, 80)),
        ("Lightning Talks", "12 speakers, 5 minutes each", (100, 70, 50)),
        ("Hackathon Kickoff", "48-hour challenge: Build an AI assistant", (50, 80, 110)),
        ("Day 2 Dinner Gala", "Awards ceremony and live entertainment", (100, 50, 70)),
        ("Day 3: Innovation Track", "Startup pitches and investor meetups", (60, 110, 90)),
        ("Startup Pitch Competition", "15 startups competed for $500K investment", (80, 100, 60)),
        ("Fireside Chat: Ethics in AI", "Dr. James Park and Prof. Lisa Wang", (50, 70, 90)),
        ("Community Meetups", "Open source project collaboration sessions", (70, 80, 100)),
        ("Closing Keynote", "CTO Robert Kim: Building What's Next", (40, 60, 110)),
        ("Event Analytics Dashboard", "Real-time engagement metrics", (80, 90, 70)),
        ("Attendee Satisfaction Survey", "Overall rating: 4.7/5.0", (60, 100, 80)),
        ("Social Media Highlights", "#TechSummit2025 trending worldwide", (100, 60, 90)),
        ("Press Coverage Summary", "Featured in 50+ tech publications", (70, 80, 60)),
        ("Sponsor ROI Report", "150% average return on sponsorship investment", (80, 70, 100)),
        ("Lessons Learned", "Key takeaways for next year's planning", (60, 80, 90)),
        ("Thank You to Volunteers", "180 volunteers made this possible", (90, 70, 80)),
        ("Save the Date: 2026 Summit", "March 20-22, 2026 | New York City", (41, 65, 122)),
    ]

    # Generate all images upfront
    print("Generating high-resolution images...")
    images = []
    for i, (title, subtitle, color) in enumerate(slide_content):
        # 3000x2000 PNG images with noise (each ~2-4MB as PNG with noise)
        img_buf = generate_large_image(3000, 2000, color, title)
        img_path = f'/tmp/event_img_{i:02d}.png'
        with open(img_path, 'wb') as f:
            f.write(img_buf.read())
        images.append(img_path)
        if (i + 1) % 10 == 0:
            print(f"  Generated {i + 1}/30 images")

    print("Building presentation...")

    # Use layout 0 (Title) for first and last slides, layout 5 (Blank) for content
    for i, (title, subtitle, color) in enumerate(slide_content):
        if i == 0 or i == 29:
            slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title layout
            if slide.shapes.title:
                slide.shapes.title.text = title
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = subtitle
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

            # Add title text box
            from pptx.util import Inches, Pt
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0]
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(41, 65, 122)

            # Add subtitle
            txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.6))
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            run2 = p2.runs[0]
            run2.font.size = Pt(18)
            run2.font.color.rgb = RGBColor(100, 100, 100)

        # Add image to every slide (background-ish or content image)
        img_path = images[i]
        if i == 0 or i == 29:
            # Full background image
            slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)
        else:
            # Content image positioned in center-bottom area
            slide.shapes.add_picture(
                img_path,
                Inches(1.5), Inches(2.0),
                Inches(10), Inches(5.0)
            )

    # Save initial version
    prs.save(OUTPUT)
    print(f"Base presentation saved: {OUTPUT}")

    # Now add 4 additional unused slide masters by manipulating the ZIP/XML directly
    # Total: 5 masters (1 used by all slides + 4 unused)
    # The task says "5 master slides defined but only 2 are actually used"
    print("Adding unused slide masters...")
    add_extra_masters(OUTPUT, count=4)

    # Clean up temp images
    for img_path in images:
        try:
            os.remove(img_path)
        except OSError:
            pass

    file_size = os.path.getsize(OUTPUT)
    print(f'Initial file created: {OUTPUT} ({file_size / (1024*1024):.1f} MB)')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=3.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


def add_extra_masters(pptx_path, count=4):
    """Add extra slide masters to the presentation by duplicating and modifying the existing master."""

    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'ct': 'http://schemas.openxmlformats.org/package/2006/content-types',
        'rel': 'http://schemas.openxmlformats.org/package/2006/relationships',
    }

    tmp_path = pptx_path + '.tmp'

    # Master names and background colors
    extra_masters = [
        ("Corporate Blue", "003366"),
        ("Presentation Dark", "1A1A2E"),
        ("Green Accent", "2D5F2D"),
        ("Warm Orange", "8B4513"),
    ][:count]

    with zipfile.ZipFile(pptx_path, 'r') as zin, zipfile.ZipFile(tmp_path, 'w') as zout:
        master1_xml = zin.read('ppt/slideMasters/slideMaster1.xml')
        master1_rels = zin.read('ppt/slideMasters/_rels/slideMaster1.xml.rels')

        pres_xml = zin.read('ppt/presentation.xml')
        pres_rels = zin.read('ppt/_rels/presentation.xml.rels')
        content_types = zin.read('[Content_Types].xml')

        pres_rels_tree = etree.fromstring(pres_rels)
        existing_rids = []
        for rel in pres_rels_tree:
            rid = rel.get('Id', '')
            if rid.startswith('rId'):
                try:
                    existing_rids.append(int(rid[3:]))
                except ValueError:
                    pass
        next_rid = max(existing_rids) + 1 if existing_rids else 10

        pres_tree = etree.fromstring(pres_xml)
        pns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        rns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

        master_id_lst = pres_tree.find(f'{{{pns}}}sldMasterIdLst')
        if master_id_lst is None:
            master_id_lst = etree.SubElement(pres_tree, f'{{{pns}}}sldMasterIdLst')

        existing_master_ids = []
        for mid in master_id_lst:
            mid_val = mid.get('id')
            if mid_val:
                existing_master_ids.append(int(mid_val))
        next_master_id = max(existing_master_ids) + 1 if existing_master_ids else 2147483649

        ct_tree = etree.fromstring(content_types)
        ct_ns = 'http://schemas.openxmlformats.org/package/2006/content-types'
        rel_type_master = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster'

        new_master_files = {}

        for idx, (master_name, bg_color) in enumerate(extra_masters):
            master_num = idx + 2
            master_path = f'ppt/slideMasters/slideMaster{master_num}.xml'
            master_rels_path = f'ppt/slideMasters/_rels/slideMaster{master_num}.xml.rels'

            master_tree = etree.fromstring(master1_xml)
            ans = 'http://schemas.openxmlformats.org/drawingml/2006/main'

            bg = master_tree.find(f'.//{{{pns}}}bg')
            if bg is None:
                cSld = master_tree.find(f'{{{pns}}}cSld')
                if cSld is not None:
                    bg = etree.SubElement(cSld, f'{{{pns}}}bg')
                    bgPr = etree.SubElement(bg, f'{{{pns}}}bgPr')
                    solidFill = etree.SubElement(bgPr, f'{{{ans}}}solidFill')
                    srgbClr = etree.SubElement(solidFill, f'{{{ans}}}srgbClr')
                    srgbClr.set('val', bg_color)
                    etree.SubElement(bgPr, f'{{{ans}}}effectLst')

            new_master_files[master_path] = etree.tostring(master_tree, xml_declaration=True, encoding='UTF-8', standalone=True)

            rels_tree = etree.fromstring(master1_rels)
            new_master_files[master_rels_path] = etree.tostring(rels_tree, xml_declaration=True, encoding='UTF-8', standalone=True)

            rid = f'rId{next_rid}'
            next_rid += 1
            master_id_elem = etree.SubElement(master_id_lst, f'{{{pns}}}sldMasterId')
            master_id_elem.set('id', str(next_master_id))
            master_id_elem.set(f'{{{rns}}}id', rid)
            next_master_id += 1

            rel_elem = etree.SubElement(pres_rels_tree, 'Relationship')
            rel_elem.set('Id', rid)
            rel_elem.set('Type', rel_type_master)
            rel_elem.set('Target', f'slideMasters/slideMaster{master_num}.xml')

            override = etree.SubElement(ct_tree, f'{{{ct_ns}}}Override')
            override.set('PartName', f'/{master_path}')
            override.set('ContentType', 'application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml')

        for item in zin.infolist():
            if item.filename == 'ppt/presentation.xml':
                zout.writestr(item, etree.tostring(pres_tree, xml_declaration=True, encoding='UTF-8', standalone=True))
            elif item.filename == 'ppt/_rels/presentation.xml.rels':
                zout.writestr(item, etree.tostring(pres_rels_tree, xml_declaration=True, encoding='UTF-8', standalone=True))
            elif item.filename == '[Content_Types].xml':
                zout.writestr(item, etree.tostring(ct_tree, xml_declaration=True, encoding='UTF-8', standalone=True))
            else:
                zout.writestr(item, zin.read(item.filename))

        for path, data in new_master_files.items():
            zout.writestr(path, data)

    shutil.move(tmp_path, pptx_path)
    print(f"Added {count} extra slide masters to {pptx_path} (total: {count + 1})")


create_initial()
