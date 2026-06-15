"""
FINAL REWARD SCRIPT - SUCCESS
Task: Right after my current Slide 2—so it becomes the new Slide 3—I want to insert one fresh page that uses the exact “Title and Two Content” layout. How do I drop that in without messing up the rest of the deck?
Generated: 2025-09-10 15:15:24
Status: success
Model: azure-o3
Total Steps: 6
"""

import os
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER


def verify_task(file_path: str) -> float:
    """Verify that a new slide (now Slide 3) was inserted with the
    exact "Title and Two Content" layout.

    Scoring (progressive):
        0.2 – Presentation has at least 3 slides (insertion point exists)
        0.5 – Slide 3 layout name matches a "Title and Two Content" pattern
        0.3 – Slide 3 contains 1 title placeholder and ≥2 content placeholders
    Returns a float score between 0.0 and 1.0 and prints detailed debug info.
    """

    print(f"Starting verification for: {file_path}")
    max_score = 1.0
    total_score = 0.0

    # ------------- Preliminary checks (no points for file existence) -------------
    if not os.path.isfile(file_path):
        print("✗ Presentation file not found.")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # ----------------------- Requirement 1: Slide count --------------------------
    slide_count = len(prs.slides)
    print(f"Slide count: {slide_count}")

    if slide_count >= 3:
        print("✓ Presentation has at least 3 slides (required for task)")
        total_score += 0.2
    else:
        print("✗ Presentation has fewer than 3 slides – cannot contain the new Slide 3")
        print("REWARD: 0.0")
        return 0.0  # Further checks make no sense

    # ------------------- Focus on the inserted slide (Slide 3) -------------------
    target_slide = prs.slides[2]  # zero-based index

    # Requirement 2: Layout name must correspond to "Title and Two Content"
    layout_name = target_slide.slide_layout.name or ""
    layout_name_lower = layout_name.lower()
    print(f"Slide 3 layout name: '{layout_name}'")

    expected_layout_keywords = [
        "title and two content",
        "title and 2 content",
        "two content",
        "title & two content",
    ]
    layout_match = any(k in layout_name_lower for k in expected_layout_keywords)

    if layout_match:
        print("✓ Slide 3 uses a layout that matches 'Title and Two Content'")
        total_score += 0.5
    else:
        print("✗ Slide 3 layout name does not indicate a 'Title and Two Content' layout")

    # Requirement 3: Placeholder structure – should include 1 title + 2 content
    title_placeholders = 0
    content_placeholders = 0

    for shape in target_slide.shapes:
        if not shape.is_placeholder:
            continue
        ph_type = shape.placeholder_format.type
        if ph_type == PP_PLACEHOLDER.TITLE:
            title_placeholders += 1
        # Content placeholders can be BODY (2) or OBJECT (7) depending on template
        if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            content_placeholders += 1

    print(
        f"Found {title_placeholders} title placeholder(s) and "
        f"{content_placeholders} content placeholder(s) on Slide 3"
    )

    if title_placeholders >= 1 and content_placeholders >= 2:
        print("✓ Slide 3 contains the correct number of placeholders")
        total_score += 0.3
    else:
        print("✗ Slide 3 placeholder structure is incorrect")

    # --------------------------- Final score output -----------------------------
    final_score = min(total_score, max_score)
    print(f"Total verification score: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score


# ------------------------------ Script execution --------------------------------
if __name__ == "__main__":
    verify_task(
        "/home/user/right_after_my_current_slide_2so_it_becomes_the_new_slide_3i_want_to_insert_one_fresh_page_that_uses_golden.pptx"
    )
