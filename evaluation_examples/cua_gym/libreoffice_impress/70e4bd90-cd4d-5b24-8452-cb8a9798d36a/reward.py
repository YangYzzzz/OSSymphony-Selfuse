"""
Reward Script: Move audio icon off visible slide area on slide 5
Task ID: impress_fix_080
Domain: libreoffice_impress
Scoring:
  Component 1 (0.6): Audio shape on slide 5 is positioned outside the visible slide area
  Component 2 (0.4): Audio shape still exists with original content (not deleted) AND is off-slide
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_080'


def persist_app_state():
    """Save any unsaved LibreOffice changes before verification."""
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0

    Task: The audio icon on slide 5 should be moved outside the visible slide area
    while keeping the audio shape functional (not deleted).

    Initial state: Audio shape "Audio Narration" at (5.00in, 3.75in) - center of slide
    Golden state: Audio shape "Audio Narration" at (-1.00in, -1.00in) - off-slide
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"PRECONDITION FAIL: Expected at least 5 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # slide 5 (0-indexed)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Find the audio narration shape on slide 5
    audio_shape = None
    for shape in slide.shapes:
        if shape.name == "Audio Narration":
            audio_shape = shape
            break

    if audio_shape is None:
        # Shape was deleted instead of moved - task requires keeping audio functional
        print("FAIL: Audio Narration shape not found on slide 5 (was it deleted?)")
        print("REWARD: 0.0")
        return 0.0

    shape_left = audio_shape.left
    shape_top = audio_shape.top
    shape_width = audio_shape.width
    shape_height = audio_shape.height

    print(f"INFO: Audio shape position: left={shape_left} ({shape_left/914400:.2f}in), "
          f"top={shape_top} ({shape_top/914400:.2f}in)")
    print(f"INFO: Audio shape size: width={shape_width} ({shape_width/914400:.2f}in), "
          f"height={shape_height} ({shape_height/914400:.2f}in)")
    print(f"INFO: Slide bounds: width={slide_width} ({slide_width/914400:.2f}in), "
          f"height={slide_height} ({slide_height/914400:.2f}in)")

    # Determine if the shape is completely outside the visible slide area.
    # The shape is off-slide if its right edge is left of slide left (0),
    # OR its bottom edge is above slide top (0),
    # OR its left edge is right of slide right edge,
    # OR its top edge is below slide bottom edge.
    shape_right = shape_left + shape_width
    shape_bottom = shape_top + shape_height

    is_off_left = shape_right <= 0  # entirely left of slide
    is_off_top = shape_bottom <= 0  # entirely above slide
    is_off_right = shape_left >= slide_width  # entirely right of slide
    is_off_bottom = shape_top >= slide_height  # entirely below slide

    is_completely_off_slide = is_off_left or is_off_top or is_off_right or is_off_bottom

    # Also check if shape is mostly off-slide (more lenient: center is off-slide)
    shape_center_x = shape_left + shape_width // 2
    shape_center_y = shape_top + shape_height // 2
    is_center_off_slide = (shape_center_x < 0 or shape_center_x > slide_width or
                           shape_center_y < 0 or shape_center_y > slide_height)

    # Component 1: Audio shape is positioned outside the visible slide area (0.6 points)
    # This checks the core task requirement: the icon must be moved off-slide.
    try:
        if is_completely_off_slide:
            print(f"PASS: Component 1 - Audio icon is completely off the visible slide area (0.6 pts)")
            total_score += 0.6
        elif is_center_off_slide:
            # Partial credit: center is off but shape partially overlaps
            print(f"PARTIAL: Component 1 - Audio icon center is off-slide but edges overlap (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 - Audio icon is still within the visible slide area "
                  f"(center at {shape_center_x/914400:.2f}in, {shape_center_y/914400:.2f}in)")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Audio shape exists with original content AND is off-slide (0.4 points)
    # This verifies the shape was moved (not deleted) and audio is still functional.
    # The shape must have text containing the music symbol and be off the slide.
    try:
        has_text = False
        if audio_shape.has_text_frame:
            shape_text = audio_shape.text_frame.text.strip()
            has_text = len(shape_text) > 0
            print(f"INFO: Audio shape text: '{shape_text}'")

        if is_completely_off_slide and has_text:
            print(f"PASS: Component 2 - Audio shape preserved with content AND off-slide (0.4 pts)")
            total_score += 0.4
        elif is_center_off_slide and has_text:
            # Partial: mostly off-slide with content
            print(f"PARTIAL: Component 2 - Audio shape has content, mostly off-slide (0.2 pts)")
            total_score += 0.2
        else:
            if not has_text:
                print(f"FAIL: Component 2 - Audio shape text is empty (content may have been removed)")
            else:
                print(f"FAIL: Component 2 - Audio shape is not off-slide")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist any unsaved GUI state before verification
persist_app_state()

# Run verification
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
