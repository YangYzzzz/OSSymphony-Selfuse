"""
Initial Setup: Move distracting audio icon off slide
Task ID: impress_fix_080
Domain: libreoffice_impress

Creates a 5-slide presentation with an embedded audio on slide 5.
The audio icon is visible at the center of the slide (distracting).
"""

import os
import shlex
import struct
import subprocess
import time
import shutil
import zipfile
import tempfile
import math
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_080'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_wav_file(filepath, duration_sec=3, sample_rate=22050):
    """Create a simple WAV file with a tone."""
    num_samples = int(duration_sec * sample_rate)
    frequency = 440  # A4 note
    amplitude = 16000

    with open(filepath, 'wb') as f:
        # WAV header
        data_size = num_samples * 2  # 16-bit mono
        f.write(b'RIFF')
        f.write(struct.pack('<I', 36 + data_size))
        f.write(b'WAVE')
        f.write(b'fmt ')
        f.write(struct.pack('<I', 16))  # chunk size
        f.write(struct.pack('<H', 1))   # PCM
        f.write(struct.pack('<H', 1))   # mono
        f.write(struct.pack('<I', sample_rate))
        f.write(struct.pack('<I', sample_rate * 2))  # byte rate
        f.write(struct.pack('<H', 2))   # block align
        f.write(struct.pack('<H', 16))  # bits per sample
        f.write(b'data')
        f.write(struct.pack('<I', data_size))

        for i in range(num_samples):
            sample = int(amplitude * math.sin(2 * math.pi * frequency * i / sample_rate))
            f.write(struct.pack('<h', sample))


def create_initial():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Quarterly Business Review"
    slide1.placeholders[1].text = "Global Operations Division — Q4 2025"

    # --- Slide 2: Revenue Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Revenue Performance"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Total Revenue: $4.2M (+12% YoY)"
    p2a = body2.add_paragraph()
    p2a.text = "APAC Region: $1.8M — strongest growth at 18%"
    p2a.level = 1
    p2b = body2.add_paragraph()
    p2b.text = "EMEA Region: $1.4M — steady at 8% growth"
    p2b.level = 1
    p2c = body2.add_paragraph()
    p2c.text = "Americas: $1.0M — slight decline of 2%"
    p2c.level = 1

    # --- Slide 3: Key Milestones ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Key Milestones Achieved"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Launched new enterprise tier with 340 sign-ups in first month"
    for milestone in [
        "Completed SOC 2 Type II certification",
        "Reduced average response time by 35%",
        "Expanded team to 120 engineers across 4 offices",
        "Strategic partnership with Meridian Analytics signed"
    ]:
        p = body3.add_paragraph()
        p.text = milestone
        p.level = 0

    # --- Slide 4: Roadmap ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "2026 Roadmap"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Q1: Launch AI-powered analytics dashboard"
    for item in [
        "Q2: Mobile app v2.0 with offline capability",
        "Q3: Enterprise SSO and advanced RBAC",
        "Q4: International expansion — LATAM and SEA markets"
    ]:
        p = body4.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 5: Summary with Audio ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Executive Summary & Narration"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "This slide includes an audio narration summarizing the key takeaways."
    p5a = body5.add_paragraph()
    p5a.text = ""
    p5b = body5.add_paragraph()
    p5b.text = "Key Points:"
    p5b.runs[0].font.bold = True
    for point in [
        "Revenue grew 12% year-over-year driven by APAC expansion",
        "Product roadmap aligned with customer feedback priorities",
        "Team scaling on track for 2026 targets"
    ]:
        p = body5.add_paragraph()
        p.text = point
        p.level = 1

    # Save the base presentation first
    tmp_pptx = f'{WORKDIR}/{TASK_ID}_tmp.pptx'
    prs.save(tmp_pptx)

    # Now create a WAV audio file
    wav_path = f'{WORKDIR}/narration.wav'
    create_wav_file(wav_path, duration_sec=3)

    # Inject audio into slide 5 by manipulating the OOXML zip
    inject_audio(tmp_pptx, OUTPUT, wav_path, slide_num=5)

    # Cleanup
    os.remove(tmp_pptx)
    os.remove(wav_path)

    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


def inject_audio(src_pptx, dst_pptx, wav_path, slide_num=5):
    """
    Inject audio into the specified slide by modifying the OOXML zip.
    Places an audio icon at the center of the slide (visible/distracting).
    """
    # Namespaces
    NS = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'rel': 'http://schemas.openxmlformats.org/package/2006/relationships',
    }

    # Read the wav file bytes
    with open(wav_path, 'rb') as f:
        wav_bytes = f.read()

    tmpdir = tempfile.mkdtemp()
    try:
        # Extract pptx
        with zipfile.ZipFile(src_pptx, 'r') as zf:
            zf.extractall(tmpdir)

        slide_xml_path = os.path.join(tmpdir, f'ppt/slides/slide{slide_num}.xml')
        rels_path = os.path.join(tmpdir, f'ppt/slides/_rels/slide{slide_num}.xml.rels')

        # Save audio file into ppt/media/
        media_dir = os.path.join(tmpdir, 'ppt/media')
        os.makedirs(media_dir, exist_ok=True)
        audio_filename = 'audio_narration.wav'
        audio_dest = os.path.join(media_dir, audio_filename)
        shutil.copy(wav_path, audio_dest)

        # Update [Content_Types].xml to include wav type
        ct_path = os.path.join(tmpdir, '[Content_Types].xml')
        ET.register_namespace('', 'http://schemas.openxmlformats.org/package/2006/content-types')
        ct_tree = ET.parse(ct_path)
        ct_root = ct_tree.getroot()
        ct_ns = 'http://schemas.openxmlformats.org/package/2006/content-types'

        # Check if wav extension already registered
        has_wav = False
        for ext_elem in ct_root.findall(f'{{{ct_ns}}}Default'):
            if ext_elem.get('Extension') == 'wav':
                has_wav = True
                break
        if not has_wav:
            wav_default = ET.SubElement(ct_root, f'{{{ct_ns}}}Default')
            wav_default.set('Extension', 'wav')
            wav_default.set('ContentType', 'audio/wav')
        ct_tree.write(ct_path, xml_declaration=True, encoding='UTF-8')

        # Add relationship in slide rels
        ET.register_namespace('', 'http://schemas.openxmlformats.org/package/2006/relationships')
        rels_tree = ET.parse(rels_path)
        rels_root = rels_tree.getroot()

        # Find next rId
        existing_ids = []
        for rel in rels_root:
            rid = rel.get('Id', '')
            if rid.startswith('rId'):
                try:
                    existing_ids.append(int(rid[3:]))
                except ValueError:
                    pass
        next_id = max(existing_ids) + 1 if existing_ids else 1
        audio_rid = f'rId{next_id}'

        # Add audio relationship
        audio_rel = ET.SubElement(rels_root, 'Relationship')
        audio_rel.set('Id', audio_rid)
        audio_rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/audio')
        audio_rel.set('Target', f'../media/{audio_filename}')

        rels_tree.write(rels_path, xml_declaration=True, encoding='UTF-8')

        # Add audio shape to slide XML
        # Position: center of slide (5in from left, 3.75in from top = center of 10x7.5 slide)
        # Using standard 10" x 7.5" slide: center at 5in, 3.75in
        # Icon size: 1in x 1in
        # Position the icon centered: left = 4.5in, top = 3.25in
        left_emu = 4572000   # ~5.0 inches (center-ish)
        top_emu = 3429000    # ~3.75 inches (center-ish)
        width_emu = 914400   # 1 inch
        height_emu = 914400  # 1 inch

        # Register all needed namespaces
        for prefix, uri in NS.items():
            if prefix != 'rel':
                ET.register_namespace(prefix, uri)
        ET.register_namespace('p', NS['p'])
        ET.register_namespace('a', NS['a'])
        ET.register_namespace('r', NS['r'])

        slide_tree = ET.parse(slide_xml_path)
        slide_root = slide_tree.getroot()

        # Find the spTree (shape tree)
        spTree = slide_root.find(f'.//{{{NS["p"]}}}cSld/{{{NS["p"]}}}spTree')
        if spTree is None:
            spTree = slide_root.find(f'.//{{{NS["p"]}}}spTree')

        # Build the audio shape XML using raw XML string for precision
        audio_shape_xml = f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="100" name="Audio Narration">
      <a:hlinkClick r:id="" action="ppaction://media"/>
    </p:cNvPr>
    <p:cNvSpPr/>
    <p:nvPr>
      <a:audioFile r:link="{audio_rid}"/>
    </p:nvPr>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{left_emu}" y="{top_emu}"/>
      <a:ext cx="{width_emu}" cy="{height_emu}"/>
    </a:xfrm>
    <a:prstGeom prst="ellipse">
      <a:avLst/>
    </a:prstGeom>
    <a:solidFill>
      <a:srgbClr val="4472C4"/>
    </a:solidFill>
    <a:ln w="12700">
      <a:solidFill>
        <a:srgbClr val="2F528F"/>
      </a:solidFill>
    </a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p>
      <a:pPr algn="ctr"/>
      <a:r>
        <a:rPr lang="en-US" sz="2400" b="1">
          <a:solidFill>
            <a:srgbClr val="FFFFFF"/>
          </a:solidFill>
        </a:rPr>
        <a:t>&#x266B;</a:t>
      </a:r>
    </a:p>
  </p:txBody>
</p:sp>'''

        audio_elem = ET.fromstring(audio_shape_xml)
        spTree.append(audio_elem)

        slide_tree.write(slide_xml_path, xml_declaration=True, encoding='UTF-8')

        # Repack into a new pptx
        with zipfile.ZipFile(dst_pptx, 'w', zipfile.ZIP_DEFLATED) as zout:
            for root_dir, dirs, files in os.walk(tmpdir):
                for fn in files:
                    full_path = os.path.join(root_dir, fn)
                    arcname = os.path.relpath(full_path, tmpdir)
                    zout.write(full_path, arcname)

    finally:
        shutil.rmtree(tmpdir)


create_initial()
