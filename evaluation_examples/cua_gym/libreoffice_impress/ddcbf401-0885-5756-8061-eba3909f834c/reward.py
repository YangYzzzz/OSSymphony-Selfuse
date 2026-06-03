"""
Reward Script: Insert chart linked to external spreadsheet data on slide 6
Task ID: impress_gf5_014
Domain: libreoffice_impress
Scoring:
  - Component 1: Chart exists on slide 6 (0.30 pts)
  - Component 2: Chart type is column/bar or line (0.15 pts)
  - Component 3: Chart categories match years 2019-2023 (0.20 pts)
  - Component 4: Chart values match revenue data from financials.xlsx (0.20 pts)
  - Component 5: Chart has external data reference (embedded/linked data source) (0.15 pts)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_014'


def persist_app_state(domain):
    """Save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide6 = prs.slides[5]  # 0-indexed

    # Find chart shape on slide 6
    chart_shape = None
    for shape in slide6.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Chart exists on slide 6 (0.30 points)
    try:
        if chart_shape is not None:
            print(f"PASS: Component 1 -- Chart found on slide 6: '{chart_shape.name}' (0.30 pts)")
            total_score += 0.30
        else:
            print("FAIL: Component 1 -- No chart found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # If no chart, remaining checks will all fail, but we still run them for diagnostics
    if chart_shape is None:
        print("FAIL: Components 2-5 skipped (no chart on slide 6)")
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    chart = chart_shape.chart

    # Component 2: Chart type is column/bar or line (0.15 points)
    try:
        from pptx.enum.chart import XL_CHART_TYPE
        chart_type = chart.chart_type
        # Accept column clustered, bar clustered, line, or similar variants
        acceptable_types = {
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            XL_CHART_TYPE.COLUMN_STACKED,
            XL_CHART_TYPE.BAR_CLUSTERED,
            XL_CHART_TYPE.BAR_STACKED,
            XL_CHART_TYPE.LINE,
            XL_CHART_TYPE.LINE_MARKERS,
        }
        if chart_type in acceptable_types:
            print(f"PASS: Component 2 -- Chart type is {chart_type} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 -- Chart type is {chart_type}, expected column/bar/line")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Chart categories match years 2019-2023 (0.20 points)
    try:
        plot = chart.plots[0]
        categories = list(plot.categories)
        expected_years = ['2019', '2020', '2021', '2022', '2023']
        # Normalize: categories might be strings or ints
        cat_strs = [str(c).strip() for c in categories]
        if cat_strs == expected_years:
            print(f"PASS: Component 3 -- Categories match years 2019-2023 (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 3 -- Categories are {cat_strs}, expected {expected_years}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Chart values match revenue data (0.20 points)
    try:
        expected_revenue = [12500000.0, 13800000.0, 16200000.0, 19500000.0, 22800000.0]
        series_values = None
        for series in chart.series:
            series_values = list(series.values)
            break  # first series

        if series_values is not None:
            # Allow small floating point tolerance
            values_match = len(series_values) == len(expected_revenue)
            if values_match:
                for actual, expected in zip(series_values, expected_revenue):
                    if actual is None or abs(actual - expected) > 1.0:
                        values_match = False
                        break

            if values_match:
                print(f"PASS: Component 4 -- Revenue values match expected data (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 4 -- Values {series_values} != expected {expected_revenue}")
        else:
            print("FAIL: Component 4 -- No series data found in chart")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Chart has external data reference (0.15 points)
    # Check for externalData element in chart XML, which indicates linked/embedded data source
    try:
        chart_xml = chart._chartSpace.xml
        has_external_data = 'externalData' in chart_xml
        if has_external_data:
            print(f"PASS: Component 5 -- Chart has external data reference (0.15 pts)")
            total_score += 0.15
        else:
            # Also check via ZIP for relationship to embedded xlsx
            import zipfile
            embedded_files = []
            with zipfile.ZipFile(file_path, 'r') as zf:
                embedded_files = [n for n in zf.namelist()
                                  if 'embed' in n.lower() and n.endswith('.xlsx')]
            if len(embedded_files) > 0:
                print(f"PASS: Component 5 -- Chart has embedded xlsx data source (0.15 pts)")
                total_score += 0.15
            else:
                print("FAIL: Component 5 -- No external data reference or embedded xlsx found")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
