"""
Initial Setup: Insert chart linked to external spreadsheet on slide 6
Task ID: impress_gf5_014
Domain: libreoffice_impress

Creates:
  1. /home/user/data/financials.xlsx - External spreadsheet with Summary sheet
  2. /home/user/impress_gf5_014.pptx - 7-slide annual report presentation
     Slide 6 is 'Financial Performance' with empty content placeholder (NO chart)
  3. Opens the presentation in LibreOffice Impress
"""

import os
import shlex
import subprocess
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_014'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
DATA_DIR = f'{WORKDIR}/data'
FINANCIALS = f'{DATA_DIR}/financials.xlsx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_financials():
    """Create the external spreadsheet with financial data."""
    import openpyxl

    os.makedirs(DATA_DIR, exist_ok=True)

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = 'Summary'

    # Headers
    ws['A1'] = 'Year'
    ws['B1'] = 'Revenue'

    # Data: 2019-2023 revenue figures (in millions)
    data = [
        (2019, 12500000),
        (2020, 13800000),
        (2021, 16200000),
        (2022, 19500000),
        (2023, 22800000),
    ]
    for i, (year, revenue) in enumerate(data, start=2):
        ws.cell(row=i, column=1, value=year)
        ws.cell(row=i, column=2, value=revenue)

    # Format revenue column as currency
    for row in range(2, 7):
        ws.cell(row=row, column=2).number_format = '$#,##0'

    # Adjust column widths
    ws.column_dimensions['A'].width = 12
    ws.column_dimensions['B'].width = 18

    # Add a second sheet with quarterly details
    ws2 = wb.create_sheet('Quarterly')
    ws2['A1'] = 'Quarter'
    ws2['B1'] = 'Revenue'
    ws2['C1'] = 'Expenses'
    ws2['D1'] = 'Net Profit'
    quarters = [
        ('Q1 2023', 5200000, 3800000, 1400000),
        ('Q2 2023', 5500000, 3900000, 1600000),
        ('Q3 2023', 5800000, 4100000, 1700000),
        ('Q4 2023', 6300000, 4300000, 2000000),
    ]
    for i, row_data in enumerate(quarters, start=2):
        for j, val in enumerate(row_data):
            ws2.cell(row=i, column=j + 1, value=val)

    wb.save(FINANCIALS)
    print(f'Financials spreadsheet created: {FINANCIALS}')


def create_presentation():
    """Create a 7-slide annual report presentation. Slide 6 has empty placeholder."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Annual Report 2023"
    slide1.placeholders[1].text = "Meridian Technologies Inc.\nFiscal Year Ended December 31, 2023"
    # Dark blue background
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x2A, 0x4A)
    # White title text
    for run in slide1.shapes.title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.size = Pt(36)
    for para in slide1.placeholders[1].text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Meridian Technologies achieved record revenue of $22.8M in 2023, representing a 17% year-over-year increase."
    p2 = body2.add_paragraph()
    p2.text = "Key highlights include expansion into three new markets, successful launch of the CloudSync platform, and a 23% increase in enterprise client acquisition."
    p2.space_before = Pt(12)
    p3 = body2.add_paragraph()
    p3.text = "Operating margin improved to 28.4%, up from 25.1% in the prior year, driven by operational efficiencies and higher-margin product mix."
    p3.space_before = Pt(12)

    # --- Slide 3: Market Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Market Overview"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "The global enterprise software market grew 12.4% in 2023 to $295B."
    items3 = [
        "Cloud infrastructure spending up 29% YoY",
        "AI/ML integration driving 40% of new enterprise deals",
        "North America remains largest market (42% share)",
        "Asia-Pacific showing fastest growth at 18.7% CAGR",
        "Cybersecurity concerns pushing 65% of firms to upgrade",
    ]
    for item in items3:
        p = body3.add_paragraph()
        p.text = item
        p.level = 1
        p.space_before = Pt(6)

    # --- Slide 4: Product Highlights ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Product Highlights"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "CloudSync Platform"
    body4.paragraphs[0].runs[0].font.bold = True
    products = [
        ("CloudSync Platform", [
            "Launched Q2 2023 with 1,200+ enterprise users by year end",
            "99.97% uptime SLA achieved across all regions",
        ]),
        ("DataVault Analytics Suite", [
            "Version 4.0 released with real-time streaming capabilities",
            "Processing 2.4 petabytes daily for top-tier clients",
        ]),
        ("SecureEdge Gateway", [
            "New zero-trust architecture adopted by 340 organizations",
            "Reduced mean breach detection time by 67%",
        ]),
    ]
    body4.clear()
    for prod_name, details in products:
        ph = body4.add_paragraph()
        ph.text = prod_name
        for run in ph.runs:
            run.font.bold = True
            run.font.size = Pt(16)
        for detail in details:
            pd = body4.add_paragraph()
            pd.text = detail
            pd.level = 1
            pd.space_before = Pt(4)

    # --- Slide 5: Team & Operations ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Team & Operations"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Our global workforce expanded to 1,847 employees across 12 offices worldwide."
    ops_items = [
        "Engineering headcount grew 34% to support CloudSync launch",
        "Employee retention rate: 91.2% (industry avg: 82%)",
        "Opened new offices in Singapore and Berlin",
        "Average revenue per employee: $123K (up from $108K)",
        "Invested $4.2M in employee development programs",
        "Diversity hiring increased 28% year-over-year",
    ]
    for item in ops_items:
        p = body5.add_paragraph()
        p.text = item
        p.level = 1
        p.space_before = Pt(6)

    # --- Slide 6: Financial Performance (EMPTY - no chart) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Financial Performance"
    # Leave the content placeholder with descriptive text but NO chart
    body6 = slide6.placeholders[1].text_frame
    body6.text = "[Chart placeholder - Revenue data from financials.xlsx to be inserted here]"
    for run in body6.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        run.font.italic = True
        run.font.size = Pt(14)

    # --- Slide 7: Outlook ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "2024 Outlook"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "Projected revenue growth of 20-25% driven by CloudSync enterprise adoption."
    outlook_items = [
        "Target: $27.5M - $28.5M in revenue",
        "Planned expansion into Latin American markets",
        "DataVault 5.0 launch scheduled for Q3 2024",
        "Strategic acquisition pipeline under evaluation",
        "R&D investment increasing to 22% of revenue",
    ]
    for item in outlook_items:
        p = body7.add_paragraph()
        p.text = item
        p.level = 1
        p.space_before = Pt(6)

    prs.save(OUTPUT)
    print(f'Presentation created: {OUTPUT}')


def main():
    create_financials()
    create_presentation()

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


main()
