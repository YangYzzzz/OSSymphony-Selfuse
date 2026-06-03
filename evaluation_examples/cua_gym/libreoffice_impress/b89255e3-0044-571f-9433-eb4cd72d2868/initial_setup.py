"""
Initial Setup: Create a syllabus presentation with 8 slides, slide 3 titled 'Grading Policy' but empty body.
Task ID: impress_teach_093
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_093'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_body_slide(prs, title_text, body_lines):
    """Add a slide with layout 1 (Title + Content), set title and body bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Introduction to Data Science"
    slide1.placeholders[1].text = "Fall 2025 — Professor Elena Martinez\nDepartment of Computer Science"

    # --- Slide 2: Course Overview ---
    add_title_body_slide(prs, "Course Overview", [
        "Explore foundations of data science and analytics",
        "Hands-on projects using Python, R, and SQL",
        "Topics: statistics, machine learning, data visualization",
        "Prerequisites: MATH 201, CS 110",
        "Class meets Tuesday & Thursday, 10:00 – 11:30 AM",
    ])

    # --- Slide 3: Grading Policy (EMPTY — task target) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add only the title as a text box at top
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Grading Policy"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

    # --- Slide 4: Required Textbooks ---
    add_title_body_slide(prs, "Required Textbooks", [
        "\"Data Science from Scratch\" by Joel Grus (2nd Edition)",
        "\"Python for Data Analysis\" by Wes McKinney",
        "\"An Introduction to Statistical Learning\" by James et al.",
        "Additional readings posted weekly on course portal",
    ])

    # --- Slide 5: Schedule Overview ---
    add_title_body_slide(prs, "Schedule Overview", [
        "Weeks 1-3: Data wrangling and exploratory analysis",
        "Weeks 4-6: Probability and statistical inference",
        "Week 7: Midterm Exam",
        "Weeks 8-11: Machine learning fundamentals",
        "Weeks 12-14: Advanced topics and guest lectures",
        "Week 15: Final project presentations",
        "Week 16: Final Exam",
    ])

    # --- Slide 6: Office Hours ---
    add_title_body_slide(prs, "Office Hours", [
        "Professor Martinez: Mon & Wed 2:00 – 4:00 PM, Room 312",
        "TA Sarah Chen: Tue & Thu 1:00 – 3:00 PM, Lab 105",
        "TA Marcus Johnson: Fri 10:00 AM – 12:00 PM, Lab 105",
        "Drop-ins welcome; appointments preferred for project reviews",
        "Email: emartinez@university.edu",
    ])

    # --- Slide 7: Academic Integrity ---
    add_title_body_slide(prs, "Academic Integrity", [
        "All submitted work must be your own",
        "Collaboration is encouraged for study, not for graded submissions",
        "Plagiarism detection tools will be used on all assignments",
        "First offense: zero on assignment; second offense: course failure",
        "Refer to University Honor Code for full policy details",
    ])

    # --- Slide 8: Questions? ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    txBox = slide8.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

    p2 = tf.add_paragraph()
    p2.text = "emartinez@university.edu"
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.size = Pt(20)
    run2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
