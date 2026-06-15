"""
Reward Script: Grading policy slide with pie chart and legend on slide 3
Task ID: impress_teach_093
Domain: libreoffice_impress
Scoring:
  Component 1: Pie chart exists on slide 3 (0.20 pts)
  Component 2: Chart has 5 correct category names (0.20 pts)
  Component 3: Chart has correct percentage values (0.20 pts)
  Component 4: Chart segments have distinct colors (0.10 pts)
  Component 5: Text box with 'Minimum 60% to pass' on slide 3 (0.15 pts)
  Component 6: That text is bold and colored #C62828 (0.15 pts)
"""

import os
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_093'


def persist_app_state(domain):
    """Try to save any unsaved LibreOffice state."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]  # 0-indexed

    # =========================================================
    # Component 1: Pie chart exists on slide 3 (0.20 points)
    # =========================================================
    chart_shape = None
    try:
        for shape in slide3.shapes:
            if hasattr(shape, 'has_chart') and shape.has_chart:
                chart_shape = shape
                break

        if chart_shape is not None:
            chart = chart_shape.chart
            # Check it's a PIE chart (chart_type == 5 for PIE)
            chart_type_val = chart.chart_type
            if chart_type_val is not None and chart_type_val == 5:
                print(f"PASS: Component 1 -- Pie chart found on slide 3 (type={chart_type_val}) (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 1 -- Chart found but type is {chart_type_val}, expected PIE (5)")
        else:
            print("FAIL: Component 1 -- No chart found on slide 3")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # =========================================================
    # Component 2: Chart has 5 correct category names (0.20 points)
    # =========================================================
    expected_categories = ['Participation', 'Homework', 'Midterm', 'Final Project', 'Final Exam']
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            plot = chart.plots[0]
            actual_cats = list(plot.categories)
            if len(actual_cats) == 5:
                # Case-insensitive comparison
                matches = sum(
                    1 for a, e in zip(actual_cats, expected_categories)
                    if a.strip().lower() == e.lower()
                )
                if matches == 5:
                    print(f"PASS: Component 2 -- All 5 categories match: {actual_cats} (0.20 pts)")
                    total_score += 0.20
                else:
                    print(f"FAIL: Component 2 -- Only {matches}/5 categories match. Actual: {actual_cats}")
            else:
                print(f"FAIL: Component 2 -- Expected 5 categories, found {len(actual_cats)}: {actual_cats}")
        else:
            print("FAIL: Component 2 -- No chart to check categories")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # =========================================================
    # Component 3: Chart has correct percentage values (0.20 points)
    # =========================================================
    expected_values = [10.0, 20.0, 25.0, 20.0, 25.0]
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            plot = chart.plots[0]
            series0 = plot.series[0]
            actual_vals = list(series0.values)
            if len(actual_vals) == 5:
                # Allow small tolerance for floating point
                val_matches = sum(
                    1 for a, e in zip(actual_vals, expected_values)
                    if abs(a - e) < 0.5
                )
                if val_matches == 5:
                    print(f"PASS: Component 3 -- All 5 values correct: {actual_vals} (0.20 pts)")
                    total_score += 0.20
                else:
                    print(f"FAIL: Component 3 -- Only {val_matches}/5 values match. Actual: {actual_vals}, Expected: {expected_values}")
            else:
                print(f"FAIL: Component 3 -- Expected 5 values, found {len(actual_vals)}: {actual_vals}")
        else:
            print("FAIL: Component 3 -- No chart to check values")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # =========================================================
    # Component 4: Chart segments have distinct colors (0.10 points)
    # =========================================================
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            ns = {
                'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
            }
            chart_xml = chart._chartSpace.xml
            root = ET.fromstring(chart_xml)
            colors = set()
            for dPt in root.findall('.//c:dPt', ns):
                srgb = dPt.find('.//a:srgbClr', ns)
                if srgb is not None:
                    colors.add(srgb.get('val'))
            if len(colors) >= 5:
                print(f"PASS: Component 4 -- {len(colors)} distinct colors found: {colors} (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 4 -- Only {len(colors)} distinct colors found (need >= 5): {colors}")
        else:
            print("FAIL: Component 4 -- No chart to check colors")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # =========================================================
    # Component 5: Text box with 'Minimum 60% to pass' on slide 3 (0.15 points)
    # =========================================================
    target_text_shape = None
    try:
        for shape in slide3.shapes:
            if shape.has_text_frame:
                full_text = shape.text_frame.text.strip()
                if 'minimum' in full_text.lower() and '60' in full_text and 'pass' in full_text.lower():
                    target_text_shape = shape
                    break

        if target_text_shape is not None:
            print(f"PASS: Component 5 -- Text box with 'Minimum 60% to pass' found: '{target_text_shape.text_frame.text.strip()}' (0.15 pts)")
            total_score += 0.15
        else:
            print("FAIL: Component 5 -- No text box containing 'Minimum 60% to pass' found on slide 3")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # =========================================================
    # Component 6: Text is bold and colored #C62828 (0.15 points)
    # =========================================================
    try:
        if target_text_shape is not None:
            bold_count = 0
            color_match_count = 0
            run_count = 0
            for para in target_text_shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        run_count += 1
                        # Check bold
                        if run.font.bold is True:
                            bold_count += 1
                        # Check color
                        try:
                            if run.font.color.type is not None:
                                rgb_str = str(run.font.color.rgb).upper()
                                if rgb_str == 'C62828':
                                    color_match_count += 1
                        except Exception:
                            pass

            if run_count > 0 and bold_count >= run_count and color_match_count >= run_count:
                print(f"PASS: Component 6 -- Text is bold ({bold_count}/{run_count} runs) and colored #C62828 ({color_match_count}/{run_count} runs) (0.15 pts)")
                total_score += 0.15
            elif bold_count > 0 and color_match_count == 0:
                print(f"PARTIAL: Component 6 -- Text is bold but color is not #C62828 (0.075 pts)")
                total_score += 0.075
            elif color_match_count > 0 and bold_count == 0:
                print(f"PARTIAL: Component 6 -- Text color is #C62828 but not bold (0.075 pts)")
                total_score += 0.075
            else:
                print(f"FAIL: Component 6 -- bold_count={bold_count}, color_match_count={color_match_count}, run_count={run_count}")
        else:
            print("FAIL: Component 6 -- No target text box to check formatting")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
