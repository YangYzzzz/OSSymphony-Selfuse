"""
Reward Script: Design professional slide master for CloudSync presentation
Task ID: impress_sales_050
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): Logo image on all 12 slides, top-left, ~0.8" size
  Component 2 (0.20): Accent line (#E94560) across bottom on all 12 slides
  Component 3 (0.15): Slide numbers bottom-right, 10pt gray (#999999) on all slides
  Component 4 (0.15): Company name 'CloudSync' bottom-left, 8pt gray (#999999) on all slides
  Component 5 (0.15): Title font Arial 28pt bold #1A1A2E on all titled slides
  Component 6 (0.15): Body font Arial 16pt #333333 on content slides
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_050'


def persist_app_state(domain):
    """Save any unsaved LibreOffice changes before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def get_font_color_rgb(run):
    """Safely get font color RGB string, returns None if not set."""
    try:
        if run.font.color.type is not None:
            return str(run.font.color.rgb)
    except Exception:
        pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    if num_slides != 12:
        print(f"WARNING: Expected 12 slides, found {num_slides}")

    # Component 1: Logo image on all slides, top-left, ~0.8" (0.20 points)
    # Golden: PICTURE shape, left=274320, top=182880, w=731520, h=731520, blob=1565 bytes
    # 0.8 inches = 731520 EMU. Tolerance: within 0.2 inches of expected position/size
    try:
        slides_with_logo = 0
        tolerance = Inches(0.2)  # position/size tolerance

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Check top-left position (left and top both near 0)
                    if (shape.left < Inches(1.0) and shape.top < Inches(1.0) and
                            abs(shape.width - Inches(0.8)) < tolerance and
                            abs(shape.height - Inches(0.8)) < tolerance):
                        slides_with_logo += 1
                        break

        logo_ratio = slides_with_logo / max(num_slides, 1)
        if logo_ratio >= 0.9:
            print(f"PASS: Component 1 -- Logo found on {slides_with_logo}/{num_slides} slides (0.20 pts)")
            total_score += 0.20
        elif logo_ratio >= 0.5:
            partial = 0.10
            print(f"PARTIAL: Component 1 -- Logo found on {slides_with_logo}/{num_slides} slides ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- Logo found on only {slides_with_logo}/{num_slides} slides")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Accent line (#E94560) across bottom on all slides (0.20 points)
    # Golden: AUTO_SHAPE, left=0, top=6400800, w=12191695 (full width), h=12700 (1pt),
    #         fill=SOLID E94560, line_width=0
    try:
        slides_with_line = 0
        slide_height = prs.slide_height  # 6858000 EMU
        bottom_zone = slide_height - Inches(1.0)  # line should be near bottom

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    # Check: near bottom, spans most of width, thin height
                    if (shape.top >= bottom_zone and
                            shape.width >= prs.slide_width * 0.8 and
                            shape.height <= Inches(0.1)):
                        # Check fill color
                        try:
                            fill = shape.fill
                            if fill.type == 1:  # SOLID
                                color = str(fill.fore_color.rgb).upper()
                                if color == 'E94560':
                                    slides_with_line += 1
                                    break
                        except Exception:
                            pass

        line_ratio = slides_with_line / max(num_slides, 1)
        if line_ratio >= 0.9:
            print(f"PASS: Component 2 -- Accent line found on {slides_with_line}/{num_slides} slides (0.20 pts)")
            total_score += 0.20
        elif line_ratio >= 0.5:
            partial = 0.10
            print(f"PARTIAL: Component 2 -- Accent line found on {slides_with_line}/{num_slides} slides ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 -- Accent line found on only {slides_with_line}/{num_slides} slides")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Slide numbers bottom-right, 10pt gray #999999 (0.15 points)
    # Golden: TextBox, left=10820095, top=6446520, text=slide_number, font=Arial 10pt #999999
    # 10pt = 127000 EMU
    try:
        slides_with_number = 0
        mid_x = prs.slide_width // 2

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == 17:  # TEXT_BOX
                    # Bottom-right: left > mid, top near bottom
                    if shape.left > mid_x and shape.top >= bottom_zone:
                        text = shape.text_frame.paragraphs[0].text.strip() if shape.text_frame.paragraphs else ""
                        if text.isdigit():
                            # Check font properties
                            runs = shape.text_frame.paragraphs[0].runs
                            if runs:
                                run = runs[0]
                                color = get_font_color_rgb(run)
                                size_ok = run.font.size is not None and abs(run.font.size - Pt(10)) < Pt(2)
                                color_ok = color is not None and color.upper() == '999999'
                                if size_ok and color_ok:
                                    slides_with_number += 1
                                    break

        num_ratio = slides_with_number / max(num_slides, 1)
        if num_ratio >= 0.9:
            print(f"PASS: Component 3 -- Slide numbers found on {slides_with_number}/{num_slides} slides (0.15 pts)")
            total_score += 0.15
        elif num_ratio >= 0.5:
            partial = 0.07
            print(f"PARTIAL: Component 3 -- Slide numbers found on {slides_with_number}/{num_slides} slides ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 -- Slide numbers found on only {slides_with_number}/{num_slides} slides")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Company name 'CloudSync' bottom-left, 8pt gray #999999 (0.15 points)
    # Golden: TextBox, left=274320, top=6446520, text='CloudSync', font=Arial 8pt #999999
    # 8pt = 101600 EMU
    try:
        slides_with_name = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 17:  # TEXT_BOX
                    # Bottom-left: left < mid, top near bottom
                    if shape.left < mid_x and shape.top >= bottom_zone:
                        text = shape.text_frame.paragraphs[0].text.strip() if shape.text_frame.paragraphs else ""
                        if 'CloudSync' in text:
                            runs = shape.text_frame.paragraphs[0].runs
                            if runs:
                                run = runs[0]
                                color = get_font_color_rgb(run)
                                size_ok = run.font.size is not None and abs(run.font.size - Pt(8)) < Pt(2)
                                color_ok = color is not None and color.upper() == '999999'
                                if size_ok and color_ok:
                                    slides_with_name += 1
                                    break

        name_ratio = slides_with_name / max(num_slides, 1)
        if name_ratio >= 0.9:
            print(f"PASS: Component 4 -- 'CloudSync' found on {slides_with_name}/{num_slides} slides (0.15 pts)")
            total_score += 0.15
        elif name_ratio >= 0.5:
            partial = 0.07
            print(f"PARTIAL: Component 4 -- 'CloudSync' found on {slides_with_name}/{num_slides} slides ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 -- 'CloudSync' found on only {slides_with_name}/{num_slides} slides")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Title font - Arial 28pt bold #1A1A2E (0.15 points)
    # Golden: Title placeholders have run.font.name='Arial', size=355600 (28pt), bold=True, color='1A1A2E'
    # 28pt = 355600 EMU
    try:
        slides_with_title_font = 0
        slides_with_title = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                try:
                    pf = shape.placeholder_format
                    if pf is None:
                        continue
                    ph_type = pf.type
                except (ValueError, AttributeError):
                    continue
                # Title placeholder types: TITLE (1), CENTER_TITLE (3)
                if ph_type in (1, 3):
                        # Check that there's actual text
                        all_runs = [r for p in shape.text_frame.paragraphs for r in p.runs if (r.text or "").strip()]
                        if all_runs:
                            slides_with_title += 1
                            # Check first run's font properties
                            run = all_runs[0]
                            name_ok = run.font.name == 'Arial'
                            size_ok = run.font.size is not None and abs(run.font.size - Pt(28)) < Pt(2)
                            bold_ok = run.font.bold is True
                            color = get_font_color_rgb(run)
                            color_ok = color is not None and color.upper() == '1A1A2E'
                            if name_ok and size_ok and bold_ok and color_ok:
                                slides_with_title_font += 1
                        break

        if slides_with_title > 0:
            title_ratio = slides_with_title_font / slides_with_title
            if title_ratio >= 0.9:
                print(f"PASS: Component 5 -- Title font correct on {slides_with_title_font}/{slides_with_title} titled slides (0.15 pts)")
                total_score += 0.15
            elif title_ratio >= 0.5:
                partial = 0.07
                print(f"PARTIAL: Component 5 -- Title font correct on {slides_with_title_font}/{slides_with_title} titled slides ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 5 -- Title font correct on only {slides_with_title_font}/{slides_with_title} titled slides")
        else:
            print(f"FAIL: Component 5 -- No title placeholders with text found")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # Component 6: Body font - Arial 16pt #333333 (0.15 points)
    # Golden: Body/content placeholders have run.font.name='Arial', size=203200 (16pt), color='333333'
    # 16pt = 203200 EMU
    try:
        slides_with_body_font = 0
        slides_with_body = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                try:
                    pf = shape.placeholder_format
                    if pf is None:
                        continue
                    ph_type = pf.type
                except (ValueError, AttributeError):
                    continue
                # Body/content/subtitle placeholder types: BODY (2), SUBTITLE (4)
                if ph_type in (2, 4):
                    all_runs = [r for p in shape.text_frame.paragraphs for r in p.runs if (r.text or "").strip()]
                    if all_runs:
                        slides_with_body += 1
                        run = all_runs[0]
                        name_ok = run.font.name == 'Arial'
                        size_ok = run.font.size is not None and abs(run.font.size - Pt(16)) < Pt(2)
                        color = get_font_color_rgb(run)
                        color_ok = color is not None and color.upper() == '333333'
                        if name_ok and size_ok and color_ok:
                            slides_with_body_font += 1
                    break

        if slides_with_body > 0:
            body_ratio = slides_with_body_font / slides_with_body
            if body_ratio >= 0.9:
                print(f"PASS: Component 6 -- Body font correct on {slides_with_body_font}/{slides_with_body} content slides (0.15 pts)")
                total_score += 0.15
            elif body_ratio >= 0.5:
                partial = 0.07
                print(f"PARTIAL: Component 6 -- Body font correct on {slides_with_body_font}/{slides_with_body} content slides ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 6 -- Body font correct on only {slides_with_body_font}/{slides_with_body} content slides")
        else:
            print(f"FAIL: Component 6 -- No body/content placeholders with text found")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
