"""
Initial Setup: Design a professional slide master for CloudSync presentation
Task ID: impress_sales_050
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_050'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
LOGO_PATH = f'{WORKDIR}/Desktop/logo.png'


def create_logo():
    """Create a simple CloudSync logo PNG on the Desktop."""
    os.makedirs(f'{WORKDIR}/Desktop', exist_ok=True)
    img = Image.new('RGBA', (200, 200), (255, 255, 255, 0))
    draw = ImageDraw.Draw(img)
    # Draw a stylized cloud shape
    draw.ellipse([30, 60, 120, 150], fill=(41, 128, 185, 255))
    draw.ellipse([70, 40, 170, 130], fill=(52, 152, 219, 255))
    draw.ellipse([100, 70, 180, 160], fill=(41, 128, 185, 255))
    draw.rectangle([50, 100, 160, 150], fill=(46, 139, 200, 255))
    # Cloud sync arrows
    draw.polygon([(90, 80), (110, 80), (100, 60)], fill=(255, 255, 255, 230))
    draw.polygon([(90, 130), (110, 130), (100, 150)], fill=(255, 255, 255, 230))
    img.save(LOGO_PATH)
    print(f'Logo created: {LOGO_PATH}')


def create_initial():
    prs = Presentation()
    # Standard widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "CloudSync Annual Sales Report"
    slide1.placeholders[1].text = "Fiscal Year 2025 Performance Overview"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Q1-Q4 Revenue Summary"
    body2.add_paragraph().text = "Regional Performance Breakdown"
    body2.add_paragraph().text = "Key Account Wins"
    body2.add_paragraph().text = "Product Line Analysis"
    body2.add_paragraph().text = "2026 Forecasts and Targets"

    # --- Slide 3: Q1 Revenue ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Q1 2025 Revenue Performance"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Total Revenue: $12.4M (+18% YoY)"
    body3.add_paragraph().text = "Enterprise Segment: $7.2M"
    body3.add_paragraph().text = "SMB Segment: $3.8M"
    body3.add_paragraph().text = "Consumer Segment: $1.4M"
    body3.add_paragraph().text = "New customer acquisition: 342 accounts"

    # --- Slide 4: Q2 Revenue ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Q2 2025 Revenue Performance"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Total Revenue: $14.1M (+22% YoY)"
    body4.add_paragraph().text = "Enterprise Segment: $8.5M"
    body4.add_paragraph().text = "SMB Segment: $4.0M"
    body4.add_paragraph().text = "Consumer Segment: $1.6M"
    body4.add_paragraph().text = "Churn rate reduced to 3.2%"

    # --- Slide 5: Q3 Revenue ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Q3 2025 Revenue Performance"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Total Revenue: $15.8M (+25% YoY)"
    body5.add_paragraph().text = "Enterprise Segment: $9.3M"
    body5.add_paragraph().text = "SMB Segment: $4.5M"
    body5.add_paragraph().text = "Consumer Segment: $2.0M"
    body5.add_paragraph().text = "Launched CloudSync Pro tier"

    # --- Slide 6: Q4 Revenue ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Q4 2025 Revenue Performance"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Total Revenue: $18.2M (+30% YoY)"
    body6.add_paragraph().text = "Enterprise Segment: $10.8M"
    body6.add_paragraph().text = "SMB Segment: $5.1M"
    body6.add_paragraph().text = "Consumer Segment: $2.3M"
    body6.add_paragraph().text = "Annual contract value up 45%"

    # --- Slide 7: Regional Breakdown ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Regional Performance"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "North America: $32.1M (53% of total)"
    body7.add_paragraph().text = "Europe: $15.7M (26%)"
    body7.add_paragraph().text = "Asia-Pacific: $9.4M (16%)"
    body7.add_paragraph().text = "Latin America: $3.3M (5%)"

    # --- Slide 8: Key Account Wins ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Key Account Wins in 2025"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Meridian Healthcare - $2.4M (3-year deal)"
    body8.add_paragraph().text = "Atlas Financial Group - $1.8M"
    body8.add_paragraph().text = "Vertex Manufacturing - $1.5M"
    body8.add_paragraph().text = "Pinnacle Retail Corp - $1.2M"
    body8.add_paragraph().text = "Nova Logistics International - $950K"

    # --- Slide 9: Product Line Analysis ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Product Line Revenue Split"
    body9 = slide9.placeholders[1].text_frame
    body9.text = "CloudSync Storage: $28.3M (47%)"
    body9.add_paragraph().text = "CloudSync Collaborate: $18.6M (31%)"
    body9.add_paragraph().text = "CloudSync Analytics: $8.4M (14%)"
    body9.add_paragraph().text = "CloudSync Pro (New): $5.2M (8%)"

    # --- Slide 10: Customer Satisfaction ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "Customer Satisfaction Metrics"
    body10 = slide10.placeholders[1].text_frame
    body10.text = "Net Promoter Score: 72 (up from 64)"
    body10.add_paragraph().text = "Customer Satisfaction: 4.6/5.0"
    body10.add_paragraph().text = "Support Response Time: <2 hours"
    body10.add_paragraph().text = "Renewal Rate: 94.8%"

    # --- Slide 11: 2026 Forecasts ---
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    slide11.shapes.title.text = "2026 Revenue Forecasts"
    body11 = slide11.placeholders[1].text_frame
    body11.text = "Projected Annual Revenue: $78M-$85M"
    body11.add_paragraph().text = "Target Growth Rate: 28-35%"
    body11.add_paragraph().text = "New Market Entry: Southeast Asia, Middle East"
    body11.add_paragraph().text = "Product Launches: CloudSync AI, CloudSync Shield"
    body11.add_paragraph().text = "Headcount Plan: 450 to 620 employees"

    # --- Slide 12: Thank You ---
    slide12 = prs.slides.add_slide(prs.slide_layouts[0])
    slide12.shapes.title.text = "Thank You"
    slide12.placeholders[1].text = "Questions? Contact: sales@cloudsync.io"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


# Execute
create_logo()
create_initial()
launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')
