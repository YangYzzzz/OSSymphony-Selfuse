"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 101 of my LibreOffice Impress file, please add a triangle that measures exactly 4 cm across the base and 3 cm in height, fill it with the “Green 3” palette color, and position it snugly in the bottom-right corner of that slide.
Generated: 2025-09-10 22:35:40
Status: success
Model: azure-o3
Total Steps: 6
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE

def verify_triangle_task(file_path: str) -> float:
    """Verify that slide 101 contains a correctly sized, colored, and positioned triangle.

    Scoring (progressive):
        0.4 – Triangle shape of correct type exists on slide 101
        0.2 – Triangle dimensions are 4 cm (width) × 3 cm (height)  (±0.1 cm tolerance)
        0.2 – Triangle fill color matches Green 3 (RGB ≈ 00A933)
        0.2 – Triangle is positioned snugly in the bottom-right corner (≤0.3 cm gap)
    Returns:
        float: total score between 0.0 and 1.0
    """

    EMU_PER_CM = 360000  # English Metric Units per centimetre
    max_score = 1.0
    score = 0.0

    # ---------- Load presentation (no points for basic loading) ----------
    if not os.path.exists(file_path):
        print('✗ File not found:', file_path)
        print('REWARD: 0.0')
        return 0.0

    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded with {len(prs.slides)} slides (no points)")
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        print('REWARD: 0.0')
        return 0.0

    # ---------- Ensure slide 101 exists ----------
    if len(prs.slides) < 101:
        print('✗ Slide 101 does not exist')
        print('REWARD: 0.0')
        return 0.0

    slide = prs.slides[100]  # zero-based index
    slide_w, slide_h = prs.slide_width, prs.slide_height

    # ---------- 1. Triangle presence (0.4) ----------
    triangle = None
    for shp in slide.shapes:
        if (
            shp.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and
            shp.auto_shape_type in (
                MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
                MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE,
            )
        ):
            triangle = shp
            break

    if triangle:
        print('✓ Triangle shape found on slide 101 (0.4)')
        score += 0.4
    else:
        print('✗ Triangle shape not found on slide 101')
        print(f'REWARD: {score}')
        return score  # cannot verify further

    # ---------- 2. Dimension check (0.2) ----------
    width_cm = triangle.width / EMU_PER_CM
    height_cm = triangle.height / EMU_PER_CM
    print(f'Triangle dimensions: {width_cm:.2f} cm × {height_cm:.2f} cm')
    if abs(width_cm - 4.0) <= 0.1 and abs(height_cm - 3.0) <= 0.1:
        print('✓ Triangle dimensions correct (0.2)')
        score += 0.2
    else:
        print('✗ Triangle dimensions incorrect (expected 4 cm × 3 cm)')

    # ---------- 3. Fill color check (0.2) ----------
    color_ok = False
    if triangle.fill and triangle.fill.type is not None:
        fc = triangle.fill.fore_color
        try:
            rgb_hex = str(fc.rgb).upper() if fc.rgb else None
            print('Triangle fill RGB:', rgb_hex)
            # Accept small variations often seen with palette mappings
            acceptable = {'00A933', '00A934', '00A932'}
            if rgb_hex and rgb_hex in acceptable:
                color_ok = True
        except Exception as e:
            print('Error reading color:', e)

    if color_ok:
        print('✓ Triangle fill color matches Green 3 (0.2)')
        score += 0.2
    else:
        print('✗ Triangle fill color does not match Green 3')

    # ---------- 4. Position check (0.2) ----------
    right_gap_cm = (slide_w - (triangle.left + triangle.width)) / EMU_PER_CM
    bottom_gap_cm = (slide_h - (triangle.top + triangle.height)) / EMU_PER_CM
    print(f'Gaps – right: {right_gap_cm:.2f} cm, bottom: {bottom_gap_cm:.2f} cm')
    if right_gap_cm <= 0.3 and bottom_gap_cm <= 0.3:
        print('✓ Triangle snug in bottom-right corner (0.2)')
        score += 0.2
    else:
        print('✗ Triangle not correctly positioned in bottom-right corner')

    # ---------- Final score ----------
    final_score = min(score, max_score)
    print(f'Total score: {final_score}')
    print(f'REWARD: {final_score}')
    return final_score


if __name__ == "__main__":
    FILE_PATH = (
        "/home/user/on_slide_101_of_my_libreoffice_impress_file_please_add_a_triangle_that_"
        "measures_exactly_4_cm_across__golden.pptx"
    )
    verify_triangle_task(FILE_PATH)
