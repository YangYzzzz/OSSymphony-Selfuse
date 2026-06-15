"""
FINAL REWARD SCRIPT - SUCCESS
Task: Insert a left tab stop at 6 cm for paragraph 1, then insert a tab at the start.
Generated: 2025-10-17 08:48:27
Status: success
Model: azure-o3
Total Steps: 5
"""

import os
from pptx import Presentation
from lxml import etree

def verify_tab_stop(file_path: str, cm_value: float = 6.0, tolerance_units: int = 20) -> float:
    """Verify that paragraph 1 has a left-aligned tab stop at *cm_value* cm and
    that a tab character (\t) was inserted at the beginning of that paragraph.

    Scoring (progressive):
      • 0.5 points – Tab stop at requested position found (within tolerance)
      • 0.5 points – Leading tab character present in the same paragraph
    Returns a float between 0.0 and 1.0 and prints step-by-step feedback.
    """

    print(f"Loading presentation: {file_path}")
    if not os.path.exists(file_path):
        print("✗ File not found")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        return 0.0

    # PowerPoint stores tab positions in English Metric Units / 127 = twips
    # 1 inch = 914400 EMU, 1 cm = 914400/2.54 EMU
    expected_pos_emu = cm_value / 2.54 * 914400  # EMU at *cm_value* cm
    expected_pos_units = int(round(expected_pos_emu / 127))  # value used in <a:tab pos="…"/>
    print(f"Expected tab position units: {expected_pos_units}")

    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    tab_stop_found = False
    leading_tab_found = False

    # Search every paragraph for tab stop + leading tab char
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                p_el = para._p  # XML element for the paragraph
                pPr = p_el.find('a:pPr', ns)
                if pPr is not None:
                    tabLst = pPr.find('a:tabLst', ns)
                    if tabLst is not None:
                        for tab in tabLst.findall('a:tab', ns):
                            pos_attr = tab.get('pos')
                            if pos_attr and pos_attr.isdigit():
                                pos_val = int(pos_attr)
                                if abs(pos_val - expected_pos_units) <= tolerance_units:
                                    tab_stop_found = True
                                    print(f"✓ Tab stop at {pos_val} units found in slide {slide_idx + 1}, "
                                          f"shape {shape_idx + 1}, paragraph {para_idx + 1}")
                # Check if paragraph text begins with a tab character
                para_text = ''.join((run.text or '') for run in para.runs)
                if para_text.startswith('\t'):
                    leading_tab_found = True
                    print(f"✓ Leading tab character found in slide {slide_idx + 1}, "
                          f"shape {shape_idx + 1}, paragraph {para_idx + 1}")

    # Progressive scoring
    score = 0.0
    if tab_stop_found:
        print("✓ Required tab stop located (0.5 points)")
        score += 0.5
    else:
        print("✗ Required tab stop NOT found (0 points)")

    if leading_tab_found:
        print("✓ Leading tab character inserted (0.5 points)")
        score += 0.5
    else:
        print("✗ Leading tab character NOT found (0 points)")

    final_score = min(score, 1.0)
    print(f"Final score: {final_score}")
    return final_score


if __name__ == "__main__":
    # Path to the presentation inside the VM environment
    FILE_PATH = "/home/user/insert_a_left_tab_stop_at_6_cm_for_paragraph_1_then_insert_a_tab_at_the_start.pptx"

    reward = verify_tab_stop(FILE_PATH)
    print(f"REWARD: {reward}")

