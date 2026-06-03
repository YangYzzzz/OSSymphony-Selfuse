"""
Reward Script: Import notes from first_aid_notes.docx into First_Aid_Training.pptx
Task ID: osworld_multi_apps_impress_notes_import_008
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): Slides 3, 4, 6, 8, 9, 10 all have non-empty notes added
  Component 2 (0.3): Notes on target slides match expected content from docx
  Component 3 (0.2): Pre-existing notes on slides 1, 2, 5, 7, 11, 12 remain unchanged
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_multi_apps_impress_notes_import_008'
PPTX_PATH = os.path.join(WORKDIR, 'First_Aid_Training.pptx')

# These slides had notes in the initial_env and must NOT be overwritten
# Slide numbers are 1-based
SLIDES_WITH_ORIGINAL_NOTES = [1, 2, 5, 7, 11, 12]
ORIGINAL_NOTES = {
    1: "Welcome participants and introduce yourself. Emphasize that first aid knowledge can make the difference between life and death. Ask participants if they have any prior first aid training to gauge the group's experience level.",
    2: "Practice the 911 call scenario as a group exercise. Stress that staying calm and speaking clearly helps dispatchers relay accurate information to responders. Remind participants that hanging up prematurely can delay help.",
    5: "Show participants how to apply a pressure bandage correctly. Demonstrate improvising a tourniquet using a belt or strip of cloth when commercial tourniquets are unavailable. Note the time of tourniquet application and communicate it to emergency responders.",
    7: "Show how to construct a basic splint using rigid material (e.g., magazines, cardboard) and bandages. Demonstrate an arm sling using a triangular bandage. Remind participants to check circulation, sensation, and movement distal to the injury before and after splinting.",
    11: "Present a local weather-related scenario relevant to participants. For heat stroke, demonstrate the ice-sheet cooling technique used by athletic trainers. Discuss the importance of buddy systems in extreme weather conditions and how to recognize early warning signs before emergencies escalate.",
    12: "End with a Q&A session. Distribute course completion certificates if applicable. Remind participants to share their training with family members. Provide handout with local emergency numbers, nearest AED locations, and information about refresher courses."
}

# These slides were empty and must have notes added from the docx
# Expected notes content for each slide that needed to be filled in
SLIDES_TO_FILL = [3, 4, 6, 8, 9, 10]
EXPECTED_NOTES_KEYWORDS = {
    3: ["chest compressions", "CPR", "sternum"],
    4: ["Heimlich", "abdominal thrusts", "obstruction"],
    6: ["rule of nines", "burn", "first-degree"],
    8: ["shock", "bleeding", "anaphylaxis"],
    9: ["EpiPen", "epinephrine", "allergen"],
    10: ["Glasgow Coma Scale", "head", "consciousness"]
}


def get_slide_notes(slide):
    """Get text from a slide's notes, stripping leading/trailing whitespace."""
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Verify slide count is correct (12 slides expected)
    if len(prs.slides) != 12:
        print(f"CRITICAL: Expected 12 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Track which target slides received notes (used by Component 3 as a gate condition)
    slides_with_notes = []

    # Component 1: All 6 target slides (3, 4, 6, 8, 9, 10) have non-empty notes (0.5 pts)
    # Each slide adds up to 0.5/6 ~= 0.0833 points
    # We award the full 0.5 only if ALL 6 target slides have notes
    try:
        slides_without_notes = []
        for slide_num in SLIDES_TO_FILL:
            slide = prs.slides[slide_num - 1]
            notes = get_slide_notes(slide)
            if notes:
                slides_with_notes.append(slide_num)
            else:
                slides_without_notes.append(slide_num)

        if len(slides_with_notes) == 6:
            print(f"PASS: Component 1 — All 6 target slides (3,4,6,8,9,10) now have non-empty notes (0.5 pts)")
            total_score += 0.5
        elif len(slides_with_notes) > 0:
            partial = round(len(slides_with_notes) / 6 * 0.5, 4)
            print(f"PARTIAL: Component 1 — {len(slides_with_notes)}/6 target slides have notes. "
                  f"Missing: {slides_without_notes}. Awarding {partial} pts")
            if partial > 0:
                total_score += partial
        else:
            print(f"FAIL: Component 1 — None of the target slides have notes. "
                  f"Missing slides: {slides_without_notes}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Notes on target slides contain expected keywords from docx content (0.3 pts)
    # Each slide earns 0.3/6 = 0.05 points for containing the expected keywords
    try:
        keyword_passes = 0
        for slide_num in SLIDES_TO_FILL:
            slide = prs.slides[slide_num - 1]
            notes = get_slide_notes(slide)
            keywords = EXPECTED_NOTES_KEYWORDS[slide_num]
            # All keywords must be present (case-insensitive) for the slide to pass
            all_keywords_found = all(kw.lower() in notes.lower() for kw in keywords)
            if all_keywords_found:
                keyword_passes += 1
                print(f"PASS: Component 2 slide {slide_num} — keywords found: {keywords}")
            else:
                missing_kw = [kw for kw in keywords if kw.lower() not in notes.lower()]
                print(f"FAIL: Component 2 slide {slide_num} — missing keywords: {missing_kw}. "
                      f"Notes preview: {notes[:60] if notes else '(empty)'}")

        keyword_score = round(keyword_passes / 6 * 0.3, 4)
        if keyword_passes == 6:
            print(f"PASS: Component 2 — All 6 target slides have correct content (0.3 pts)")
            total_score += keyword_score
        elif keyword_passes > 0:
            print(f"PARTIAL: Component 2 — {keyword_passes}/6 slides have correct content ({keyword_score} pts)")
            total_score += keyword_score
        else:
            print(f"FAIL: Component 2 — No target slides have correct content (0.0 pts)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Original notes on slides 1, 2, 5, 7, 11, 12 remain unchanged (0.2 pts)
    # This component ONLY awards points if:
    #   (a) the target slides (3, 4, 6, 8, 9, 10) have notes (the task was performed), AND
    #   (b) the pre-existing notes were not overwritten
    # This makes the component conditional on task completion — it FAILS on initial_env
    # because we gate on slides_with_notes having been populated.
    try:
        # Gate: the task must have been performed (at least some target slides have notes)
        if len(slides_with_notes) == 0:
            print("FAIL: Component 3 — Task not performed (no target slides have notes); "
                  "skipping preservation check")
        else:
            preserved_correctly = []
            overwritten_slides = []
            for slide_num in SLIDES_WITH_ORIGINAL_NOTES:
                slide = prs.slides[slide_num - 1]
                actual_notes = get_slide_notes(slide)
                expected_notes = ORIGINAL_NOTES[slide_num]
                if actual_notes == expected_notes:
                    preserved_correctly.append(slide_num)
                else:
                    overwritten_slides.append(slide_num)
                    print(f"FAIL: Component 3 slide {slide_num} — notes changed. "
                          f"Expected start: {repr(expected_notes[:50])}, "
                          f"Actual start: {repr(actual_notes[:50] if actual_notes else '(empty)')}")

            if len(preserved_correctly) == 6:
                print(f"PASS: Component 3 — All 6 original notes preserved unchanged (0.2 pts)")
                total_score += 0.2
            elif len(preserved_correctly) > 0:
                partial = round(len(preserved_correctly) / 6 * 0.2, 4)
                print(f"PARTIAL: Component 3 — {len(preserved_correctly)}/6 original notes preserved. "
                      f"Overwritten: {overwritten_slides}. Awarding {partial} pts")
                if partial > 0:
                    total_score += partial
            else:
                print(f"FAIL: Component 3 — Original notes were overwritten. Slides affected: {overwritten_slides}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 4)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Run verification on the canonical artifact path
if not os.path.exists(PPTX_PATH):
    print(f"File not found: {PPTX_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(PPTX_PATH)
