"""
Initial Setup: First Aid Training presentation with partial presenter notes
Task ID: osworld_multi_apps_impress_notes_import_008
Domain: libreoffice_impress

Creates:
  - /home/user/First_Aid_Training.pptx  (12 slides; slides 1,2,5,7,11,12 have notes)
  - /home/user/Desktop/first_aid_notes.docx  (notes for all 12 slides)
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document

WORKDIR = '/home/user'
TASK_ID = 'osworld_multi_apps_impress_notes_import_008'
PPTX_OUTPUT = f'{WORKDIR}/First_Aid_Training.pptx'
DESKTOP_DIR = f'{WORKDIR}/Desktop'
DOCX_OUTPUT = f'{DESKTOP_DIR}/first_aid_notes.docx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


# --- Slide content definitions ---
SLIDES = [
    {
        "title": "Introduction to First Aid",
        "body": (
            "• Definition and importance of first aid\n"
            "• The 3 Ps: Preserve life, Prevent deterioration, Promote recovery\n"
            "• Legal considerations: duty of care and Good Samaritan laws\n"
            "• Basic first aid kit contents and maintenance"
        ),
    },
    {
        "title": "Calling for Emergency Help",
        "body": (
            "• When to call emergency services (911 / local number)\n"
            "• Information to provide: location, nature of emergency, number of casualties\n"
            "• Staying on the line with the dispatcher\n"
            "• Directing bystanders to meet responders"
        ),
    },
    {
        "title": "CPR Basics",
        "body": (
            "• Recognize cardiac arrest: unresponsive, not breathing normally\n"
            "• Call 911 immediately and request an AED\n"
            "• Chest compressions: 30 compressions at 100-120 bpm, 2 inches deep\n"
            "• Rescue breaths: 2 breaths after every 30 compressions (if trained)\n"
            "• Continue until AED arrives or professional help takes over"
        ),
    },
    {
        "title": "Choking Response",
        "body": (
            "• Universal choking sign: hands clutched to throat\n"
            "• Encourage coughing if victim can cough or speak\n"
            "• Abdominal thrusts (Heimlich maneuver): 5 back blows + 5 abdominal thrusts\n"
            "• For infants: 5 back blows + 5 chest thrusts\n"
            "• Unconscious victim: begin CPR and check mouth before rescue breaths"
        ),
    },
    {
        "title": "Wound Care and Bleeding Control",
        "body": (
            "• Apply firm direct pressure using a clean cloth or bandage\n"
            "• Elevate the injured limb above heart level when possible\n"
            "• Do not remove embedded objects — stabilize in place\n"
            "• Apply tourniquet 2-3 inches above wound for life-threatening limb bleeding\n"
            "• Monitor for signs of shock: pale skin, rapid weak pulse, confusion"
        ),
    },
    {
        "title": "Burns Treatment",
        "body": (
            "• Remove victim from heat source; ensure scene safety\n"
            "• Cool the burn with cool (not cold) running water for 10-20 minutes\n"
            "• Do NOT apply ice, butter, or home remedies\n"
            "• Cover with a sterile non-stick dressing or cling film\n"
            "• Seek medical attention for burns larger than the victim's palm"
        ),
    },
    {
        "title": "Fractures and Sprains",
        "body": (
            "• Signs of fracture: pain, swelling, deformity, inability to bear weight\n"
            "• Immobilize the injury using a splint or sling — do not straighten\n"
            "• Apply ice pack wrapped in cloth to reduce swelling\n"
            "• RICE for sprains: Rest, Ice, Compression, Elevation\n"
            "• Seek X-ray for suspected fractures; do not delay for severe injuries"
        ),
    },
    {
        "title": "Shock Recognition and Treatment",
        "body": (
            "• Types: hypovolemic, anaphylactic, cardiogenic, septic\n"
            "• Signs: pale/clammy skin, rapid weak pulse, confusion, shallow breathing\n"
            "• Lay victim flat and elevate legs 8-12 inches (unless head/spinal injury)\n"
            "• Keep victim warm with a blanket; do not give food or drink\n"
            "• Monitor breathing and pulse continuously until help arrives"
        ),
    },
    {
        "title": "Allergic Reactions and Anaphylaxis",
        "body": (
            "• Mild reaction: hives, itching, localized swelling — antihistamine if available\n"
            "• Anaphylaxis signs: throat swelling, difficulty breathing, drop in blood pressure\n"
            "• Administer epinephrine auto-injector (EpiPen) into outer thigh\n"
            "• Call 911 immediately; second dose may be needed after 5-15 minutes\n"
            "• Position: sitting up if breathing difficulty; lying flat if in shock"
        ),
    },
    {
        "title": "Head Injuries",
        "body": (
            "• Assess consciousness using AVPU: Alert, Voice, Pain, Unresponsive\n"
            "• Do not remove helmets unless victim's airway is compromised\n"
            "• Keep head and neck still if spinal injury is suspected\n"
            "• Monitor for concussion signs: headache, confusion, vomiting, unequal pupils\n"
            "• Any loss of consciousness requires immediate emergency medical evaluation"
        ),
    },
    {
        "title": "Heat and Cold Emergencies",
        "body": (
            "• Heat exhaustion: move to cool area, hydrate with water/sports drink, apply cool cloths\n"
            "• Heat stroke (temp >104°F): call 911, rapid cooling with ice packs to neck/armpits/groin\n"
            "• Hypothermia: move to warm environment, remove wet clothing, warm core first\n"
            "• Frostbite: rewarm in warm water 99-102°F; do not rub or walk on frostbitten feet\n"
            "• Prevent: appropriate clothing, stay hydrated, avoid prolonged exposure"
        ),
    },
    {
        "title": "Course Summary and Review",
        "body": (
            "• First aid saves lives when delivered promptly and correctly\n"
            "• Key skills: CPR, bleeding control, choking response, burn care, shock management\n"
            "• Always call professional services — first aid is a bridge to definitive care\n"
            "• Practice skills regularly; certification renewal every 2 years recommended\n"
            "• Resources: Red Cross, American Heart Association, local training centers"
        ),
    },
]

# Presenter notes for all 12 slides
ALL_NOTES = [
    # Slide 1 — already in pptx (initial state)
    "Welcome participants and introduce yourself. Emphasize that first aid knowledge can make the difference between life and death. Ask participants if they have any prior first aid training to gauge the group's experience level.",
    # Slide 2 — already in pptx (initial state)
    "Practice the 911 call scenario as a group exercise. Stress that staying calm and speaking clearly helps dispatchers relay accurate information to responders. Remind participants that hanging up prematurely can delay help.",
    # Slide 3 — MISSING in initial pptx (to be added by agent)
    "Demonstrate the correct hand placement for chest compressions: heel of dominant hand on lower half of sternum, interlock fingers, keep arms straight. Use a CPR manikin if available. Emphasize compression depth and rate. Note that hands-only CPR is acceptable for untrained bystanders.",
    # Slide 4 — MISSING in initial pptx (to be added by agent)
    "Invite a volunteer to demonstrate the Heimlich maneuver on a partner. Remind participants to stand slightly behind and to the side of the victim. For pregnant or obese victims, chest thrusts replace abdominal thrusts. Emphasize calling 911 if the obstruction is not cleared after several cycles.",
    # Slide 5 — already in pptx (initial state)
    "Show participants how to apply a pressure bandage correctly. Demonstrate improvising a tourniquet using a belt or strip of cloth when commercial tourniquets are unavailable. Note the time of tourniquet application and communicate it to emergency responders.",
    # Slide 6 — MISSING in initial pptx (to be added by agent)
    "Use the rule of nines to help participants estimate burn surface area: head/neck 9%, each arm 9%, chest 18%, back 18%, each leg 18%, genitalia 1%. Clarify the difference between superficial (first-degree), partial-thickness (second-degree), and full-thickness (third-degree) burns and appropriate responses for each.",
    # Slide 7 — already in pptx (initial state)
    "Show how to construct a basic splint using rigid material (e.g., magazines, cardboard) and bandages. Demonstrate an arm sling using a triangular bandage. Remind participants to check circulation, sensation, and movement distal to the injury before and after splinting.",
    # Slide 8 — MISSING in initial pptx (to be added by agent)
    "Discuss the most common causes of shock seen in everyday emergencies: severe bleeding, severe burns, and anaphylaxis. Reinforce that treating the underlying cause is critical alongside supportive positioning. Ask participants to identify shock signs in a brief case study scenario.",
    # Slide 9 — MISSING in initial pptx (to be added by agent)
    "If your facility or community has EpiPen-trained users, walk through the steps: remove safety cap, press firmly against outer thigh, hold for 10 seconds, then rub the site. Discuss that epinephrine buys time and 911 is still mandatory even if symptoms improve. Mention common allergen triggers: peanuts, tree nuts, shellfish, bee stings.",
    # Slide 10 — MISSING in initial pptx (to be added by agent)
    "Use the Glasgow Coma Scale as a brief reference: Eye opening (1-4), Verbal response (1-5), Motor response (1-6). A score below 13 warrants urgent care. Remind participants not to give anything by mouth to a head-injured person and to document any changes in consciousness for paramedics.",
    # Slide 11 — already in pptx (initial state)
    "Present a local weather-related scenario relevant to participants. For heat stroke, demonstrate the ice-sheet cooling technique used by athletic trainers. Discuss the importance of buddy systems in extreme weather conditions and how to recognize early warning signs before emergencies escalate.",
    # Slide 12 — already in pptx (initial state)
    "End with a Q&A session. Distribute course completion certificates if applicable. Remind participants to share their training with family members. Provide handout with local emergency numbers, nearest AED locations, and information about refresher courses.",
]

# Slides that already have notes in the initial pptx (1-indexed)
SLIDES_WITH_INITIAL_NOTES = {1, 2, 5, 7, 11, 12}


def create_pptx():
    """Create First_Aid_Training.pptx with 12 slides; only some slides have notes."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]   # Title Slide
    content_layout = prs.slide_layouts[1]  # Title + Content

    for idx, slide_info in enumerate(SLIDES):
        slide_num = idx + 1

        if slide_num == 1:
            slide = prs.slides.add_slide(title_layout)
            slide.shapes.title.text = slide_info["title"]
            slide.placeholders[1].text = "Emergency Response Training Program\nCertified First Aid Course"
        else:
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = slide_info["title"]
            tf = slide.placeholders[1].text_frame
            tf.text = slide_info["body"]

        # Add notes ONLY for slides that already have them in initial state
        if slide_num in SLIDES_WITH_INITIAL_NOTES:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_tf.text = ALL_NOTES[idx]

    prs.save(PPTX_OUTPUT)
    print(f'Initial PPTX created: {PPTX_OUTPUT}')


def create_docx():
    """Create first_aid_notes.docx on Desktop with notes for all 12 slides."""
    os.makedirs(DESKTOP_DIR, exist_ok=True)

    doc = Document()
    doc.add_heading('First Aid Training — Presenter Notes', level=0)

    for idx, slide_info in enumerate(SLIDES):
        slide_num = idx + 1
        doc.add_heading(f'Slide {slide_num}: {slide_info["title"]}', level=1)
        doc.add_paragraph(ALL_NOTES[idx])

    doc.save(DOCX_OUTPUT)
    print(f'Notes DOCX created: {DOCX_OUTPUT}')


def main():
    create_pptx()
    create_docx()

    # GUI-ready startup: open First_Aid_Training.pptx in LibreOffice Impress
    # and first_aid_notes.docx in LibreOffice Writer
    launch_gui(f'libreoffice --impress "{PPTX_OUTPUT}"', delay_sec=3.0)
    launch_gui(f'libreoffice --writer "{DOCX_OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress and Writer with DISPLAY=:0')


main()
