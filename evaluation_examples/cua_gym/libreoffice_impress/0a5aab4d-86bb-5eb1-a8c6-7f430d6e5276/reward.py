"""
Reward Script: Customer Case Study Slides
Task ID: impress_sales_048
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25) - Presentation has exactly 10 slides (4 new case study slides added)
  Component 2 (0.30) - Each of slides 7-10 has the correct company name title
  Component 3 (0.25) - Each case study slide has Challenge and Solution text areas
  Component 4 (0.20) - Each case study slide has a results metric section at the bottom
"""

import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_048'

# Expected company names for slides 7-10 (0-indexed: 6-9)
EXPECTED_COMPANIES = ['Acme Corp', 'TechStart', 'GlobalBank', 'MediHealth']


def get_all_text_shapes(slide):
    """Get all shapes with text frames, including nested in groups."""
    def extract(shape):
        results = []
        if hasattr(shape, "text") and hasattr(shape, "text_frame"):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.extend(extract(sub))
        return results
    out = []
    for shape in slide.shapes:
        out.extend(extract(shape))
    return out


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Presentation has exactly 10 slides (0.25 points)
    # Initial has 6 slides; golden should have 10 (6 original + 4 case studies)
    try:
        if num_slides == 10:
            print(f"PASS: Component 1 — Slide count is 10 (0.25 pts)")
            total_score += 0.25
        elif num_slides > 6:
            # Partial: some slides added but not exactly 4
            partial = min(0.15, 0.15 * (num_slides - 6) / 4)
            print(f"PARTIAL: Component 1 — Expected 10 slides, found {num_slides} (+{partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — Expected 10 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Gate: Need at least 7 slides to check case study content
    if num_slides <= 6:
        print("FAIL: No case study slides found (only original 6 slides)")
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: Each of the new slides has the correct company name as title (0.30 points)
    # 0.075 per correct company name
    try:
        company_score = 0.0
        case_slides = list(prs.slides)[6:10]  # slides 7-10 (0-indexed: 6-9)
        for idx, (slide, expected_company) in enumerate(zip(case_slides, EXPECTED_COMPANIES)):
            text_shapes = get_all_text_shapes(slide)
            # Look for the company name in any text shape on this slide
            found_company = False
            for shape in text_shapes:
                full_text = shape.text_frame.text.strip()
                if expected_company.lower() in full_text.lower():
                    found_company = True
                    break
            if found_company:
                print(f"PASS: Component 2.{idx+1} — Slide {idx+7} has company name '{expected_company}' (0.075 pts)")
                company_score += 0.075
            else:
                all_texts = [s.text_frame.text.strip()[:50] for s in text_shapes if s.text_frame.text.strip()]
                print(f"FAIL: Component 2.{idx+1} — Slide {idx+7} missing '{expected_company}'. Found texts: {all_texts}")
        total_score += company_score
        print(f"  Component 2 subtotal: {company_score}/0.30")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Each case study slide has Challenge and Solution text areas (0.25 points)
    # 0.0625 per slide (checks both Challenge AND Solution present)
    try:
        cs_score = 0.0
        case_slides = list(prs.slides)[6:10]
        for idx, slide in enumerate(case_slides):
            text_shapes = get_all_text_shapes(slide)
            has_challenge = False
            has_solution = False
            for shape in text_shapes:
                text = shape.text_frame.text.strip().lower()
                if 'challenge' in text:
                    has_challenge = True
                if 'solution' in text:
                    has_solution = True

            if has_challenge and has_solution:
                print(f"PASS: Component 3.{idx+1} — Slide {idx+7} has both Challenge and Solution areas (0.0625 pts)")
                cs_score += 0.0625
            else:
                print(f"FAIL: Component 3.{idx+1} — Slide {idx+7}: Challenge={'found' if has_challenge else 'MISSING'}, Solution={'found' if has_solution else 'MISSING'}")
        total_score += cs_score
        print(f"  Component 3 subtotal: {cs_score}/0.25")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Each case study slide has a results metric section at bottom (0.20 points)
    # We check for a text shape in the lower portion of the slide that contains metric-like content
    # (numbers, percentages, dollar amounts)
    try:
        results_score = 0.0
        case_slides = list(prs.slides)[6:10]
        slide_height = prs.slide_height
        bottom_threshold = slide_height * 0.65  # results should be in bottom 35%

        for idx, slide in enumerate(case_slides):
            text_shapes = get_all_text_shapes(slide)
            has_bottom_metrics = False
            for shape in text_shapes:
                text = shape.text_frame.text.strip()
                # Check if shape is in bottom portion and contains metric-like content
                if shape.top >= bottom_threshold and len(text) > 5:
                    # Look for numbers, percentages, or dollar signs as metric indicators
                    import re
                    if re.search(r'(\d+%|\$[\d,.]+|\d+\.\d+)', text):
                        has_bottom_metrics = True
                        break

            if has_bottom_metrics:
                print(f"PASS: Component 4.{idx+1} — Slide {idx+7} has results metrics at bottom (0.05 pts)")
                results_score += 0.05
            else:
                print(f"FAIL: Component 4.{idx+1} — Slide {idx+7} missing results metrics in bottom section")
        total_score += results_score
        print(f"  Component 4 subtotal: {results_score}/0.20")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
