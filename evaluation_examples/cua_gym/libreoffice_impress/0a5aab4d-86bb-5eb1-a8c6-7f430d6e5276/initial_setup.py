"""
Initial Setup: Create a 6-slide sales pitch presentation for case studies
Task ID: impress_sales_048
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_048'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a formatted text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    # Standard 16:9 dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Customer Case Studies"
    slide1.placeholders[1].text = "Q4 2025 Success Stories"

    # ---- Slide 2: Agenda ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_textbox(slide2, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Agenda", font_size=32, bold=True,
                color=RGBColor(0x1F, 0x49, 0x7D))
    agenda_items = [
        "1. Executive Summary",
        "2. Market Landscape",
        "3. Product Overview",
        "4. Implementation Timeline",
        "5. Customer Case Studies",
        "6. Next Steps"
    ]
    txBox = slide2.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(agenda_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.space_after = Pt(12)
        run = p.runs[0]
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ---- Slide 3: Executive Summary ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide3, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Executive Summary", font_size=32, bold=True,
                color=RGBColor(0x1F, 0x49, 0x7D))
    summary_text = (
        "Our enterprise solutions have delivered measurable ROI across "
        "diverse industries. In Q4 2025, we onboarded 47 new enterprise "
        "clients, achieving a 98.2% customer satisfaction score. Our platform "
        "processed over 12 million transactions with 99.97% uptime, "
        "demonstrating reliability at scale."
    )
    add_textbox(slide3, Inches(1.0), Inches(2.0), Inches(11), Inches(4),
                summary_text, font_size=18, color=RGBColor(0x44, 0x44, 0x44))

    # ---- Slide 4: Market Landscape ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide4, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Market Landscape", font_size=32, bold=True,
                color=RGBColor(0x1F, 0x49, 0x7D))
    market_items = [
        "Total addressable market: $84.2B (2025 est.)",
        "Year-over-year growth: 23.5%",
        "Enterprise segment accounts for 62% of revenue",
        "Key verticals: Financial Services, Healthcare, Technology, Manufacturing",
        "Competitive advantage: AI-driven automation reduces deployment time by 40%"
    ]
    txBox = slide4.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(market_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.space_after = Pt(14)
        run = p.runs[0]
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # ---- Slide 5: Product Overview ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide5, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Product Overview", font_size=32, bold=True,
                color=RGBColor(0x1F, 0x49, 0x7D))
    features = [
        "Real-time Analytics Dashboard - Monitor KPIs across all business units",
        "Automated Workflow Engine - Reduce manual processes by up to 75%",
        "Enterprise Security Suite - SOC 2 Type II and ISO 27001 certified",
        "Integration Hub - 200+ pre-built connectors for major platforms",
        "Custom Reporting - Drag-and-drop report builder with scheduled delivery"
    ]
    txBox = slide5.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(features):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.space_after = Pt(14)
        run = p.runs[0]
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # ---- Slide 6: Implementation Timeline ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide6, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Implementation Timeline", font_size=32, bold=True,
                color=RGBColor(0x1F, 0x49, 0x7D))
    timeline = [
        "Week 1-2: Discovery & Requirements Gathering",
        "Week 3-4: Platform Configuration & Data Migration",
        "Week 5-6: Integration Testing & UAT",
        "Week 7-8: Training & Go-Live Support",
        "Ongoing: Quarterly Business Reviews & Optimization"
    ]
    txBox = slide6.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(timeline):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.space_after = Pt(14)
        run = p.runs[0]
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
