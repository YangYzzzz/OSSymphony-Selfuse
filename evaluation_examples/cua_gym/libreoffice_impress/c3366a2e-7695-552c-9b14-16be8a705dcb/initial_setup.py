"""
Initial Setup: Create a 10-slide business presentation for accessibility export task
Task ID: impstruct_049
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'accessible_deck'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q3 2025 Strategic Planning Review"
    slide1.placeholders[1].text = "Prepared by the Office of Corporate Strategy\nAugust 2025"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    tf2 = slide2.placeholders[1].text_frame
    tf2.clear()
    items = [
        "Executive Summary & Key Highlights",
        "Revenue Performance by Region",
        "Product Portfolio Update",
        "Customer Acquisition Metrics",
        "Operational Efficiency Initiatives",
        "Risk Assessment & Mitigation",
        "Talent & Workforce Planning",
        "Technology Roadmap",
        "Budget Allocation for Q4",
    ]
    for i, item in enumerate(items):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = item
        p.space_after = Pt(6)

    # --- Slide 3: Executive Summary ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Executive Summary"
    tf3 = slide3.placeholders[1].text_frame
    tf3.clear()
    summary_items = [
        "Total revenue reached $128.4M, up 14.2% YoY",
        "APAC region exceeded target by 8.7%, driven by enterprise deals",
        "Customer retention rate improved to 94.3% from 91.1%",
        "Three new product lines launched ahead of schedule",
        "Operating margin expanded 220 basis points to 18.6%",
    ]
    for i, item in enumerate(summary_items):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p.text = item
        p.space_after = Pt(4)

    # --- Slide 4: Revenue by Region (table) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    txBox.text_frame.paragraphs[0].text = "Revenue Performance by Region"
    txBox.text_frame.paragraphs[0].font.size = Pt(28)
    txBox.text_frame.paragraphs[0].font.bold = True

    tbl_shape = slide4.shapes.add_table(6, 4, Inches(0.8), Inches(1.4), Inches(10), Inches(4))
    tbl = tbl_shape.table
    headers = ["Region", "Q3 Revenue ($M)", "Target ($M)", "Variance (%)"]
    data = [
        ["North America", "52.3", "50.0", "+4.6%"],
        ["EMEA", "34.1", "33.5", "+1.8%"],
        ["APAC", "28.7", "26.4", "+8.7%"],
        ["Latin America", "8.9", "9.0", "-1.1%"],
        ["Middle East & Africa", "4.4", "4.1", "+7.3%"],
    ]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row in enumerate(data, 1):
        for c, val in enumerate(row):
            tbl.cell(r, c).text = val

    # --- Slide 5: Product Portfolio ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Product Portfolio Update"
    tf5 = slide5.placeholders[1].text_frame
    tf5.clear()
    products = [
        "CloudSync Enterprise: 12,400 active licenses (+23% QoQ)",
        "DataVault Pro: Successfully migrated 340 enterprise clients",
        "InsightIQ Analytics: Beta launch with 89 pilot customers",
        "SecureEdge Gateway: FedRAMP certification obtained in July",
        "FlowBuilder 2.0: Released with 47 new workflow templates",
    ]
    for i, item in enumerate(products):
        p = tf5.paragraphs[0] if i == 0 else tf5.add_paragraph()
        p.text = item
        p.space_after = Pt(4)

    # --- Slide 6: Customer Metrics ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Customer Acquisition & Retention"
    tf6 = slide6.placeholders[1].text_frame
    tf6.clear()
    metrics = [
        "New enterprise clients signed: 87 (target: 75)",
        "Average deal size increased to $234K from $198K",
        "Net Promoter Score: 72 (industry benchmark: 58)",
        "Churn rate decreased to 5.7% from 8.9% YoY",
        "Customer lifetime value grew 31% to $1.24M",
        "Support ticket resolution time: 4.2 hours (down from 6.8)",
    ]
    for i, item in enumerate(metrics):
        p = tf6.paragraphs[0] if i == 0 else tf6.add_paragraph()
        p.text = item
        p.space_after = Pt(4)

    # --- Slide 7: Operational Efficiency ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Operational Efficiency Initiatives"
    tf7 = slide7.placeholders[1].text_frame
    tf7.clear()
    ops = [
        "Automated CI/CD pipeline reduced deployment time by 62%",
        "Cloud infrastructure costs optimized: saved $2.1M annually",
        "Vendor consolidation completed for 14 SaaS tools",
        "ISO 27001 recertification achieved with zero findings",
        "Cross-functional sprint teams reduced time-to-market by 3 weeks",
    ]
    for i, item in enumerate(ops):
        p = tf7.paragraphs[0] if i == 0 else tf7.add_paragraph()
        p.text = item
        p.space_after = Pt(4)

    # --- Slide 8: Risk Assessment ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    txBox8.text_frame.paragraphs[0].text = "Risk Assessment & Mitigation"
    txBox8.text_frame.paragraphs[0].font.size = Pt(28)
    txBox8.text_frame.paragraphs[0].font.bold = True

    tbl8_shape = slide8.shapes.add_table(5, 3, Inches(0.8), Inches(1.4), Inches(10), Inches(4))
    tbl8 = tbl8_shape.table
    risk_headers = ["Risk Factor", "Severity", "Mitigation Plan"]
    risk_data = [
        ["Supply chain disruption", "High", "Diversified to 3 secondary suppliers"],
        ["Regulatory compliance (EU AI Act)", "Medium", "Legal review and product audit in progress"],
        ["Key personnel attrition", "Medium", "Retention packages and succession planning"],
        ["Currency fluctuation impact", "Low", "Hedging strategy covers 80% of exposure"],
    ]
    for c, h in enumerate(risk_headers):
        cell = tbl8.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row in enumerate(risk_data, 1):
        for c, val in enumerate(row):
            tbl8.cell(r, c).text = val

    # --- Slide 9: Technology Roadmap ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Technology Roadmap - H2 2025"
    tf9 = slide9.placeholders[1].text_frame
    tf9.clear()
    roadmap = [
        "September: Launch AI-powered document classification module",
        "October: Complete Kubernetes migration for all microservices",
        "November: Release mobile SDK v3.0 with offline capabilities",
        "December: Deploy real-time analytics dashboard for enterprise tier",
        "Q1 2026: Begin integration with SAP S/4HANA and Salesforce",
    ]
    for i, item in enumerate(roadmap):
        p = tf9.paragraphs[0] if i == 0 else tf9.add_paragraph()
        p.text = item
        p.space_after = Pt(4)

    # --- Slide 10: Next Steps ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "Next Steps & Action Items"
    tf10 = slide10.placeholders[1].text_frame
    tf10.clear()
    next_steps = [
        "Finalize Q4 budget allocations by September 5th",
        "Schedule deep-dive sessions with each product team",
        "Present customer expansion strategy to the board on Sept 15",
        "Complete hiring plan for 45 new engineering positions",
        "Begin RFP process for data center expansion in Frankfurt",
        "Review and approve updated security policies by Sept 30",
    ]
    for i, item in enumerate(next_steps):
        p = tf10.paragraphs[0] if i == 0 else tf10.add_paragraph()
        p.text = item
        p.space_after = Pt(4)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open presentation in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=3.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
