"""
Reward Script: Before-and-after comparison presentation
Task ID: impress_wf_004
Domain: libreoffice_impress
Scoring:
  Component 1 (0.2): Renovation.pptx exists on Desktop with exactly 3 slides
  Component 2 (0.2): Slide 1 title is 'Office Renovation'
  Component 3 (0.3): Slides 2-3 each have 'Before' and 'After' text labels
  Component 4 (0.3): Slides 2-3 each have 2 images matching before/after source blobs
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_004'
DESKTOP = os.path.join(WORKDIR, 'Desktop')


def get_all_text(slide):
    """Recursively extract all text from slide shapes (including groups)."""
    texts = []
    def extract(shape):
        if hasattr(shape, 'text') and hasattr(shape, 'text_frame'):
            texts.append(shape.text.strip())
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                extract(sub)
    for shape in slide.shapes:
        extract(shape)
    return texts


def get_image_blobs(slide):
    """Get all image blobs from a slide."""
    blobs = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            blobs.append(shape.image.blob)
    return blobs


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Load source image blobs for comparison
    before_blob = None
    after_blob = None
    try:
        with open(os.path.join(DESKTOP, 'before_office.jpg'), 'rb') as f:
            before_blob = f.read()
        with open(os.path.join(DESKTOP, 'after_office.jpg'), 'rb') as f:
            after_blob = f.read()
    except Exception as e:
        print(f"WARNING: Could not load source images for comparison: {e}")

    # Component 1: File has exactly 3 slides (0.2 points)
    try:
        num_slides = len(prs.slides)
        if num_slides == 3:
            print(f"PASS: Component 1 - Presentation has 3 slides (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 1 - Expected 3 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Slide 1 title is 'Office Renovation' (0.2 points)
    try:
        if len(prs.slides) >= 1:
            slide1 = prs.slides[0]
            all_text = get_all_text(slide1)
            # Check if any text shape contains 'Office Renovation'
            title_found = any('office renovation' in t.lower() for t in all_text)
            if title_found:
                print(f"PASS: Component 2 - Slide 1 contains 'Office Renovation' title (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 2 - Slide 1 texts: {all_text}, none match 'Office Renovation'")
        else:
            print(f"FAIL: Component 2 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3a: Slide 2 has 'Before' and 'After' text labels (0.15 points)
    try:
        if len(prs.slides) >= 2:
            slide2 = prs.slides[1]
            s2_text = [t.lower() for t in get_all_text(slide2)]
            if any('before' in t for t in s2_text) and any('after' in t for t in s2_text):
                print(f"PASS: Component 3a - Slide 2 has 'Before' and 'After' labels (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 3a - Slide 2 texts: {s2_text}")
        else:
            print(f"FAIL: Component 3a - Slide 2 does not exist")
    except Exception as e:
        print(f"ERROR: Component 3a - {e}")

    # Component 3b: Slide 3 has 'Before' and 'After' text labels (0.15 points)
    try:
        if len(prs.slides) >= 3:
            slide3 = prs.slides[2]
            s3_text = [t.lower() for t in get_all_text(slide3)]
            if any('before' in t for t in s3_text) and any('after' in t for t in s3_text):
                print(f"PASS: Component 3b - Slide 3 has 'Before' and 'After' labels (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 3b - Slide 3 texts: {s3_text}")
        else:
            print(f"FAIL: Component 3b - Slide 3 does not exist")
    except Exception as e:
        print(f"ERROR: Component 3b - {e}")

    # Component 4a: Slide 2 has 2 images matching before/after source blobs (0.15 points)
    try:
        if len(prs.slides) >= 2:
            slide2 = prs.slides[1]
            img_blobs_s2 = get_image_blobs(slide2)
            if len(img_blobs_s2) >= 2:
                has_before_s2 = before_blob is not None and any(b == before_blob for b in img_blobs_s2)
                has_after_s2 = after_blob is not None and any(b == after_blob for b in img_blobs_s2)
                if has_before_s2 and has_after_s2:
                    print(f"PASS: Component 4a - Slide 2 has both before and after images (0.15 pts)")
                    total_score += 0.15
                elif len(img_blobs_s2) >= 2:
                    print(f"PARTIAL: Component 4a - Slide 2 has {len(img_blobs_s2)} images but blob mismatch (0.10 pts)")
                    total_score += 0.10
            else:
                print(f"FAIL: Component 4a - Slide 2 has {len(img_blobs_s2)} images, expected 2")
        else:
            print(f"FAIL: Component 4a - Slide 2 does not exist")
    except Exception as e:
        print(f"ERROR: Component 4a - {e}")

    # Component 4b: Slide 3 has 2 images matching before/after source blobs (0.15 points)
    try:
        if len(prs.slides) >= 3:
            slide3 = prs.slides[2]
            img_blobs_s3 = get_image_blobs(slide3)
            if len(img_blobs_s3) >= 2:
                has_before_s3 = before_blob is not None and any(b == before_blob for b in img_blobs_s3)
                has_after_s3 = after_blob is not None and any(b == after_blob for b in img_blobs_s3)
                if has_before_s3 and has_after_s3:
                    print(f"PASS: Component 4b - Slide 3 has both before and after images (0.15 pts)")
                    total_score += 0.15
                elif len(img_blobs_s3) >= 2:
                    print(f"PARTIAL: Component 4b - Slide 3 has {len(img_blobs_s3)} images but blob mismatch (0.10 pts)")
                    total_score += 0.10
            else:
                print(f"FAIL: Component 4b - Slide 3 has {len(img_blobs_s3)} images, expected 2")
        else:
            print(f"FAIL: Component 4b - Slide 3 does not exist")
    except Exception as e:
        print(f"ERROR: Component 4b - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = os.path.join(DESKTOP, 'Renovation.pptx')
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
