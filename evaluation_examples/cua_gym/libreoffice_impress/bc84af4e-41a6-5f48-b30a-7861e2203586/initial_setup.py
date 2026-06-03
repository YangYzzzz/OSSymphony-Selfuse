"""
Initial Setup: Create an 8-slide sales review presentation with slide 6 titled 'Sales Performance' (title only, no chart).
Task ID: impress_tm_063
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_063'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def add_title_only_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
    return slide


def create_initial():
    prs = Presentation()

    # Slide 1: Title slide
    add_title_slide(prs, "Annual Sales Review 2025", "Acme Global Corporation")

    # Slide 2: Agenda
    add_content_slide(prs, "Agenda", [
        "Company Overview and Market Position",
        "Regional Sales Performance",
        "Product Category Breakdown",
        "Quarterly Sales Analysis",
        "Customer Acquisition Metrics",
        "2026 Projections and Strategy",
    ])

    # Slide 3: Company Overview
    add_content_slide(prs, "Company Overview", [
        "Total Revenue: $14.2M (up 12% YoY)",
        "Active Clients: 387 across 28 countries",
        "Employee Count: 245 (up from 198 in 2024)",
        "New Market Entries: Southeast Asia, Eastern Europe",
        "Customer Retention Rate: 94.3%",
    ])

    # Slide 4: Regional Breakdown
    add_content_slide(prs, "Regional Sales Breakdown", [
        "North America: $6.8M (48% of total)",
        "Europe: $3.9M (27% of total)",
        "Asia-Pacific: $2.1M (15% of total)",
        "Latin America: $0.9M (6% of total)",
        "Middle East & Africa: $0.5M (4% of total)",
    ])

    # Slide 5: Product Categories
    add_content_slide(prs, "Product Category Performance", [
        "Enterprise Solutions: $5.4M (+18%)",
        "Cloud Services: $4.1M (+32%)",
        "Professional Services: $2.8M (+8%)",
        "Support & Maintenance: $1.9M (+5%)",
    ])

    # Slide 6: Sales Performance - TITLE ONLY, NO CHART
    add_title_only_slide(prs, "Sales Performance")

    # Slide 7: Customer Acquisition
    add_content_slide(prs, "Customer Acquisition", [
        "New Customers in 2025: 63",
        "Average Deal Size: $78,400",
        "Sales Cycle: 42 days (reduced from 58 in 2024)",
        "Conversion Rate: 28.4%",
        "Top Channel: Direct Sales (52%)",
    ])

    # Slide 8: 2026 Outlook
    add_content_slide(prs, "2026 Outlook & Strategy", [
        "Revenue Target: $16.5M (+16%)",
        "Focus Areas: AI-powered products, APAC expansion",
        "Hiring Plan: 35 new positions across engineering and sales",
        "Key Initiative: Launch self-service platform by Q2",
        "Customer Retention Goal: 96%",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
