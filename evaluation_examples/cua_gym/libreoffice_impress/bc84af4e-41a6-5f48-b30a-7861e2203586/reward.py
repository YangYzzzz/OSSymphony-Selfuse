"""
Reward Script: Insert a bar chart on slide 6 with quarterly sales data
Task ID: impress_tm_063
Domain: libreoffice_impress
Scoring:
  - Component 1: Chart exists on slide 6 (0.20)
  - Component 2: Chart is bar/column type (0.15)
  - Component 3: Categories match Q1-Q4 (0.25)
  - Component 4: Values match 120K, 145K, 98K, 167K (0.25)
  - Component 5: Chart has a title (0.15)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_063'


def persist_app_state(domain: str):
    """Attempt to save any unsaved LibreOffice state."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify that slide 6 has a bar chart with the correct quarterly sales data.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"PRECONDITION FAIL: Expected at least 6 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[5]  # slide 6, 0-indexed

    # Find chart shape on slide 6
    chart_shape = None
    for shape in slide.shapes:
        if hasattr(shape, 'chart'):
            chart_shape = shape
            break

    # Component 1: Chart exists on slide 6 (0.20 points)
    try:
        if chart_shape is not None:
            print(f"PASS: Component 1 -- Chart found on slide 6 (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 -- No chart found on slide 6")
            # No chart means nothing else can pass
            print(f"\nScore: {total_score}/1.0")
            print(f"REWARD: {total_score}")
            return total_score
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")
        print(f"REWARD: {total_score}")
        return total_score

    chart = chart_shape.chart

    # Component 2: Chart is bar/column type (0.15 points)
    # In PowerPoint, "bar chart" typically means COLUMN_CLUSTERED (vertical bars)
    # or BAR_CLUSTERED (horizontal bars). Both are acceptable for "bar chart".
    try:
        ct = chart.chart_type
        # BAR types: 57 (BAR_CLUSTERED), 58 (BAR_STACKED), 59, 60
        # COLUMN types: 51 (COLUMN_CLUSTERED), 52 (COLUMN_STACKED), 53, 54
        bar_column_types = [51, 52, 53, 54, 57, 58, 59, 60]
        if ct in bar_column_types:
            print(f"PASS: Component 2 -- Chart type is bar/column ({ct}) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 -- Expected bar/column chart type, found {ct}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Categories are Q1, Q2, Q3, Q4 (0.25 points)
    try:
        plot = chart.plots[0]
        categories = list(plot.categories)
        expected_cats = ['Q1', 'Q2', 'Q3', 'Q4']
        # Normalize: strip whitespace, case-insensitive comparison
        norm_cats = [str(c).strip() for c in categories]
        norm_expected = [str(c).strip() for c in expected_cats]

        if len(norm_cats) == len(norm_expected):
            matches = sum(1 for a, b in zip(norm_cats, norm_expected)
                          if a.upper() == b.upper())
            if matches == 4:
                print(f"PASS: Component 3 -- All 4 categories match: {categories} (0.25 pts)")
                total_score += 0.25
            else:
                # Partial: give proportional credit
                partial = 0.25 * (matches / 4)
                print(f"PARTIAL: Component 3 -- {matches}/4 categories match: {categories} ({partial:.2f} pts)")
                total_score += partial
        else:
            print(f"FAIL: Component 3 -- Expected 4 categories, found {len(categories)}: {categories}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Values match 120000, 145000, 98000, 167000 (0.25 points)
    try:
        series = chart.series[0]
        values = list(series.values)
        expected_vals = [120000.0, 145000.0, 98000.0, 167000.0]

        if len(values) == len(expected_vals):
            matches = 0
            for actual, expected in zip(values, expected_vals):
                # Allow small tolerance for floating point
                if actual is not None and abs(float(actual) - expected) < 1.0:
                    matches += 1
            if matches == 4:
                print(f"PASS: Component 4 -- All 4 values match: {values} (0.25 pts)")
                total_score += 0.25
            else:
                partial = 0.25 * (matches / 4)
                print(f"PARTIAL: Component 4 -- {matches}/4 values match: {values} ({partial:.2f} pts)")
                total_score += partial
        else:
            print(f"FAIL: Component 4 -- Expected 4 values, found {len(values)}: {values}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Chart has a title (0.15 points)
    try:
        if chart.has_title:
            title_text = chart.chart_title.text_frame.text.strip()
            if title_text:
                print(f"PASS: Component 5 -- Chart has title: '{title_text}' (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 5 -- Chart title is empty")
        else:
            print(f"FAIL: Component 5 -- Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
