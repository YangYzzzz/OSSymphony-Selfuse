"""
Initial Setup: Investor Briefing presentation with 10 slides, no timings or transitions
Task ID: impress_gf4_038
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_038'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.paragraphs[0].text = point
        else:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Investor Briefing Q2 2025",
        "Meridian Technologies Inc.\nConfidential — For Authorized Investors Only"
    )

    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "Revenue grew 34% year-over-year to $127.3M in Q2 2025",
        "Operating margin expanded to 18.2%, up from 14.7% in Q2 2024",
        "Customer base increased to 2,840 enterprise accounts (+22% YoY)",
        "Net Promoter Score reached 72, highest in company history",
        "Free cash flow of $31.5M, representing 24.7% FCF margin",
    ])

    # Slide 3: Market Overview
    add_content_slide(prs, "Market Overview", [
        "Total addressable market estimated at $48B by 2027 (Gartner)",
        "Cloud infrastructure spending accelerated 28% across industries",
        "Enterprise AI adoption rate reached 67% among Fortune 500 companies",
        "Regulatory tailwinds from EU Digital Markets Act driving demand",
        "Three new geographic markets entered: Brazil, South Korea, UAE",
    ])

    # Slide 4: Revenue Performance
    add_content_slide(prs, "Revenue Performance", [
        "Subscription revenue: $98.6M (+41% YoY) — 77% of total revenue",
        "Professional services: $18.2M (+12% YoY) — transitioning to partner-led",
        "Annual Recurring Revenue (ARR): $412M, up from $298M a year ago",
        "Net revenue retention rate: 128%, driven by platform upsell",
        "Average deal size increased 19% to $145K annually",
    ])

    # Slide 5: Financial Charts & Analysis
    add_content_slide(prs, "Financial Analysis — Detailed Breakdown", [
        "Gross margin: 74.3% (vs 71.8% in Q2 2024) — infrastructure optimization",
        "R&D investment: $28.4M (22.3% of revenue) — AI/ML capabilities",
        "Sales & marketing: $33.7M (26.5%) — efficiency improving as brand strengthens",
        "G&A: $12.1M (9.5%) — scaled admin costs with automation",
        "EBITDA: $29.8M (23.4% margin) — path to 30% by 2026",
        "Capital expenditure: $8.2M — data center expansion in Frankfurt and Singapore",
    ])

    # Slide 6: Product Roadmap
    add_content_slide(prs, "Product Roadmap", [
        "Q3 2025: Launch AI-powered analytics dashboard (Project Aurora)",
        "Q4 2025: Multi-cloud orchestration engine — AWS, Azure, GCP unified",
        "Q1 2026: Real-time collaboration suite with integrated video",
        "Q2 2026: Advanced threat detection module (FedRAMP certification pending)",
        "Ongoing: API ecosystem expansion — 340+ integrations by year-end",
    ])

    # Slide 7: Customer Acquisition
    add_content_slide(prs, "Customer Acquisition & Retention", [
        "New logos: 187 enterprise customers added in Q2 2025",
        "Win rate improved to 38% from 31% in competitive deals",
        "Customer acquisition cost (CAC) decreased 14% to $42,300",
        "Payback period reduced to 11 months from 15 months",
        "Top new clients: Siemens, Rakuten, Petrobras, Woolworths Group",
    ])

    # Slide 8: Competitive Landscape
    add_content_slide(prs, "Competitive Landscape", [
        "Positioned as Leader in Forrester Wave for Cloud Management (Q1 2025)",
        "Key differentiator: unified observability + automated remediation",
        "Primary competitors: Datadog, Dynatrace, Splunk (Cisco)",
        "Win rate vs. Datadog: 52% (up from 44% in 2024)",
        "Patent portfolio expanded to 89 granted patents, 34 pending",
    ])

    # Slide 9: Growth Strategy
    add_content_slide(prs, "Growth Strategy 2025–2027", [
        "Land-and-expand: target 150% net dollar retention by Q4 2026",
        "International expansion: APAC and LATAM to reach 35% of revenue",
        "Strategic M&A: $200M allocated for complementary acquisitions",
        "Partner ecosystem: 120+ certified solution partners by 2026",
        "Talent: hiring 400+ engineers across Bangalore, Berlin, and Austin offices",
    ])

    # Slide 10: Q&A / Contact
    add_content_slide(prs, "Questions & Contact Information", [
        "Investor Relations: Victoria Langford — ir@meridiantech.com",
        "CEO: Dr. Nathan Prescott — nathan.prescott@meridiantech.com",
        "CFO: Rachel Okonkwo — rachel.okonkwo@meridiantech.com",
        "Next Earnings Call: October 14, 2025 at 4:30 PM ET",
        "Annual Investor Day: November 18, 2025 — San Francisco, CA",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
