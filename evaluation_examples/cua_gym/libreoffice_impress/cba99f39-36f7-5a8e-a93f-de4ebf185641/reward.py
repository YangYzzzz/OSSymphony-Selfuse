"""
Reward Script: Move the data table on slide 4 to the bottom third of the slide.
Task ID: osworld_impress_table_position_bottom_003
Domain: libreoffice_impress
Scoring:
  Gate (not scored): Table found on slide 4, data integrity verified (5 rows x 4 cols, header intact)
  Component 1 (0.6): Table top position is in the bottom third of the slide (top >= slide_height * 2/3 = 5.0in)
  Component 2 (0.4): Table top is within 5%% of the golden reference position (5.0 inches / 4572000 EMU),
                     AND table data integrity preserved (5 rows x 4 cols, header='Product Line')

The two-part Component 2 ensures we score only the deliberate repositioning to the golden target position
while confirming the table content was not corrupted in the process. Data integrity alone is not scored
(it is a precondition in both initial and golden), but it must be true together with the precise position.
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_table_position_bottom_003'


def persist_app_state():
    """Send Ctrl+S to save any unsaved LibreOffice Impress edits."""
    try:
        import time
        os.environ["DISPLAY"] = ":0"
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def is_approximately_equal(val1, val2, tolerance=0.05):
    """Check if two values are within tolerance fraction of each other."""
    if val1 == val2:
        return True
    if val1 == 0 or val2 == 0:
        return val1 == val2
    return abs(val1 - val2) / max(abs(val1), abs(val2)) <= tolerance


def verify_task(file_path):
    """
    Verify that the table on slide 4 has been moved to the bottom third of the slide.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Check slide count
    if len(prs.slides) < 4:
        print(f"CRITICAL: Expected at least 4 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide4 = prs.slides[3]  # 0-indexed, slide 4
    slide_height = prs.slide_height  # EMU

    # Bottom third threshold: top must be >= slide_height * 2/3
    # For 7.5-inch slide: 6858000 * 2/3 = 4572000 EMU = 5.0 inches
    bottom_third_threshold = int(slide_height * 2 / 3)  # 4572000 EMU

    # Find the table on slide 4
    table_shape = None
    for shape in slide4.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_shape = shape
            break

    if table_shape is None:
        print("CRITICAL: No table found on slide 4 -- cannot score")
        print("REWARD: 0.0")
        return 0.0

    table_top = table_shape.top  # EMU
    table_top_in = table_top / 914400
    threshold_in = bottom_third_threshold / 914400

    # Component 1: Table top is in the bottom third (top >= slide_height * 2/3)
    # Initial state: top=3017520 EMU (3.3 in) -- NOT in bottom third -- FAILS
    # Golden state:  top=4572000 EMU (5.0 in) -- at bottom third boundary -- PASSES
    try:
        if table_top >= bottom_third_threshold:
            print(f"PASS: Component 1 -- Table is in bottom third: top={table_top_in:.4f}in "
                  f">= threshold={threshold_in:.4f}in (0.6 pts)")
            total_score += 0.6
        else:
            print(f"FAIL: Component 1 -- Table NOT in bottom third: top={table_top_in:.4f}in "
                  f"< threshold={threshold_in:.4f}in (expected top >= {threshold_in:.4f}in)")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Table top is approximately at the golden reference position
    #              AND table data integrity preserved (verifies move was done correctly, not a re-creation)
    # Golden reference: 4572000 EMU (5.0 inches, exactly at the bottom-third boundary)
    # Tolerance: 5% relative
    # Initial state: top=3017520 EMU (3.3 in) -- position check FAILS -- overall FAILS
    # Golden state:  top=4572000 EMU (5.0 in), 5 rows x 4 cols, header='Product Line' -- PASSES
    GOLDEN_TOP_EMU = 4572000  # 5.0 inches
    try:
        t = table_shape.table
        num_rows = len(t.rows)
        num_cols = len(t.columns)
        header_text = t.cell(0, 0).text.strip() if num_rows > 0 and num_cols > 0 else ""
        position_ok = is_approximately_equal(table_top, GOLDEN_TOP_EMU, tolerance=0.05)
        data_ok = (num_rows == 5 and num_cols == 4 and header_text == "Product Line")

        if position_ok and data_ok:
            print(f"PASS: Component 2 -- Table at golden position ({table_top_in:.4f}in ~= {GOLDEN_TOP_EMU/914400:.4f}in) AND data intact ({num_rows}x{num_cols}, header='{header_text}') (0.4 pts)")
            total_score += 0.4
        elif not position_ok:
            deviation_pct = abs(table_top - GOLDEN_TOP_EMU) / max(abs(table_top), abs(GOLDEN_TOP_EMU)) * 100
            print(f"FAIL: Component 2 -- Table position deviates from golden: "
                  f"{table_top_in:.4f}in vs expected ~{GOLDEN_TOP_EMU/914400:.4f}in "
                  f"(deviation: {deviation_pct:.1f}%)")
        else:
            print(f"FAIL: Component 2 -- Table data integrity failed: "
                  f"rows={num_rows} (expected 5), cols={num_cols} (expected 4), "
                  f"header='{header_text}' (expected 'Product Line')")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point: test against canonical artifact path in a given env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
persist_app_state()
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
