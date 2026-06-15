"""
Initial Setup: 5-slide analytics deck with table on slide 4 in CENTER position
Task ID: osworld_impress_table_position_bottom_003
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_table_position_bottom_003'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard widescreen: 10 inches wide, 7.5 inches tall
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_width = prs.slide_width   # 9144000 EMU
    slide_height = prs.slide_height  # 6858000 EMU

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q1 2025 Business Analytics Report"
    slide1.placeholders[1].text = "Performance Overview & Key Metrics"

    # ---- Slide 2: Executive Summary ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Revenue grew 18% year-over-year"
    p2 = tf2.add_paragraph()
    p2.text = "Customer acquisition cost down by 12%"
    p3 = tf2.add_paragraph()
    p3.text = "Net Promoter Score reached all-time high of 72"
    p4 = tf2.add_paragraph()
    p4.text = "EBITDA margin improved to 24.3%"

    # ---- Slide 3: Regional Performance ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Regional Performance"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "North America: $4.2M (+22%)"
    p3b = tf3.add_paragraph()
    p3b.text = "Europe: $2.8M (+15%)"
    p3c = tf3.add_paragraph()
    p3c.text = "APAC: $1.9M (+31%)"
    p3d = tf3.add_paragraph()
    p3d.text = "Latin America: $0.7M (+9%)"

    # ---- Slide 4: Data Analysis (chart upper area + table in CENTER) ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Title text box at top
    title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9.0), Inches(0.7))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "Q1 Product Line Performance"
    p_title.alignment = PP_ALIGN.CENTER
    run_title = p_title.runs[0]
    run_title.font.size = Pt(28)
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Chart placeholder area (upper portion, using a text box as a chart label)
    chart_label = slide4.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.4))
    tf_chart = chart_label.text_frame
    p_chart = tf_chart.paragraphs[0]
    p_chart.text = "Monthly Revenue Trend (Jan-Mar 2025)"
    p_chart.alignment = PP_ALIGN.CENTER
    run_chart = p_chart.runs[0]
    run_chart.font.size = Pt(14)
    run_chart.font.italic = True
    run_chart.font.color.rgb = RGBColor(0x70, 0x70, 0x70)

    # Simulated chart area (colored rectangle as chart background)
    from pptx.util import Emu as EmuU
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    chart_area = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(2.5))
    chart_fill = chart_area.fill
    chart_fill.solid()
    chart_fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
    tf_ca = chart_area.text_frame
    p_ca = tf_ca.paragraphs[0]
    p_ca.text = "[Chart: Monthly Revenue Bar Chart]"
    p_ca.alignment = PP_ALIGN.CENTER
    run_ca = p_ca.runs[0]
    run_ca.font.size = Pt(12)
    run_ca.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

    # ---- TABLE in CENTER of slide (initial position) ----
    # Center of slide: roughly rows 4.2 to 7.5 inches vertically = 3.5" from top
    # Table at center: top = ~3.3 inches from top
    table_rows = 5
    table_cols = 4
    table_left = Inches(0.5)
    table_top = Inches(3.3)   # CENTER of slide - this is what must be moved
    table_width = Inches(9.0)
    table_height = Inches(2.5)

    table_shape = slide4.shapes.add_table(
        table_rows, table_cols,
        table_left, table_top,
        table_width, table_height
    )
    table = table_shape.table

    # Set column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(2.5)
    table.columns[3].width = Inches(2.0)

    # Header row
    headers = ["Product Line", "Units Sold", "Revenue ($K)", "Growth %"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.size = Pt(13)
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            para.alignment = PP_ALIGN.CENTER
        # Header background
        from pptx.oxml.ns import qn as qname
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        solidFill = etree.SubElement(tcPr, qname('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qname('a:srgbClr'))
        srgbClr.set('val', '1F497D')

    # Data rows
    data = [
        ["Enterprise Software", "1,240", "$892.4", "+24.3%"],
        ["Cloud Services", "3,875", "$1,243.7", "+31.8%"],
        ["Hardware Solutions", "642", "$478.2", "+8.5%"],
        ["Professional Services", "289", "$347.6", "+18.2%"],
    ]
    row_colors = ["FFFFFF", "F2F7FF", "FFFFFF", "F2F7FF"]
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
                para.alignment = PP_ALIGN.CENTER
            # Alternate row background
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            solidFill = etree.SubElement(tcPr, qname('a:solidFill'))
            srgbClr = etree.SubElement(solidFill, qname('a:srgbClr'))
            srgbClr.set('val', row_colors[r - 1])

    # ---- Slide 5: Conclusions ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Conclusions & Next Steps"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Accelerate Cloud Services expansion in APAC"
    p5b = tf5.add_paragraph()
    p5b.text = "Invest in Enterprise Software sales team (Q2 target: +30%)"
    p5c = tf5.add_paragraph()
    p5c.text = "Optimize Hardware margin through supply chain improvements"
    p5d = tf5.add_paragraph()
    p5d.text = "Launch Professional Services certification program"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
