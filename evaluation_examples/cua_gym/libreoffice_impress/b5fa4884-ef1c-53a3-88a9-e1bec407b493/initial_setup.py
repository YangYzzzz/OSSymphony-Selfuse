"""
Initial Setup: Create a 6-slide Quarterly Review presentation.
Slide 2 has 'Key Metrics' title and 5 bullet points. No animations.
Task ID: impress_ma_053
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_053'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Quarterly Business Review"
    slide.placeholders[1].text = "Q1 2025 — Prepared by the Strategy Team"
    return slide


def add_key_metrics_slide(prs):
    """Slide 2: Key Metrics with 5 bullet points (NO animations)."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide.shapes.title.text = "Key Metrics"

    # Bulleted list in content placeholder
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    bullets = [
        "Total revenue reached $12.4M, a 15% increase over Q4 2024",
        "Customer acquisition cost decreased to $38 per new user",
        "Net promoter score improved from 62 to 71 across all regions",
        "Monthly active users grew by 23% to 1.8 million",
        "Employee retention rate held steady at 94% company-wide",
    ]

    for i, text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(18)

    return slide


def add_revenue_breakdown_slide(prs):
    """Slide 3: Revenue Breakdown."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Revenue Breakdown"

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    items = [
        "Enterprise Contracts: $5.2M (42%)",
        "SaaS Subscriptions: $4.1M (33%)",
        "Professional Services: $1.9M (15%)",
        "Licensing & Partnerships: $1.2M (10%)",
    ]
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(18)

    return slide


def add_team_highlights_slide(prs):
    """Slide 4: Team Highlights."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Team Highlights"

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    items = [
        "Engineering shipped 3 major product releases ahead of schedule",
        "Marketing launched the 'Elevate' campaign with 2.3M impressions",
        "Sales closed 47 new enterprise accounts in APAC region",
        "Customer Success reduced average ticket resolution time by 18%",
    ]
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(18)

    return slide


def add_challenges_slide(prs):
    """Slide 5: Challenges & Risks."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Challenges & Risks"

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    items = [
        "Supply chain delays impacted hardware rollout timelines",
        "Increased competition in the mid-market SaaS segment",
        "Regulatory changes in EU require compliance updates by Q3",
    ]
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(18)

    return slide


def add_next_steps_slide(prs):
    """Slide 6: Next Steps."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Next Steps & Q2 Priorities"

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    items = [
        "Finalize product roadmap for Q2 with cross-functional alignment",
        "Expand sales presence in Latin American markets",
        "Launch customer loyalty program targeting top 500 accounts",
        "Complete SOC 2 Type II audit by end of May 2025",
    ]
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(18)

    return slide


def create_initial():
    prs = Presentation()

    add_title_slide(prs)
    add_key_metrics_slide(prs)
    add_revenue_breakdown_slide(prs)
    add_team_highlights_slide(prs)
    add_challenges_slide(prs)
    add_next_steps_slide(prs)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
