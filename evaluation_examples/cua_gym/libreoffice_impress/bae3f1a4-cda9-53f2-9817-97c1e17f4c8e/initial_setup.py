"""
Initial Setup: Create a 9-slide Sales Analysis presentation with slide 6 having
title 'Quarterly Sales Breakdown' and an empty content area.
Task ID: impress_gf2_012
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_012'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullets):
    """Add a slide with Title + Content layout (layout index 1)."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, text in enumerate(bullets):
        if i == 0:
            tf.paragraphs[0].text = text
        else:
            p = tf.add_paragraph()
            p.text = text
    return slide


def add_title_only_slide(prs, title_text):
    """Add a slide with Title Only layout (layout index 5)."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title_text
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    add_title_slide(prs, "Sales Analysis 2025",
                    "Annual Performance Review\nPrepared by the Analytics Team")

    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "Total revenue reached $4.2M across all product lines",
        "Year-over-year growth of 18% exceeded projections",
        "ProductA maintained market leadership with 45% share",
        "New market expansion contributed 12% of total revenue",
        "Customer retention rate improved to 94%",
    ])

    # Slide 3: Revenue Overview
    add_content_slide(prs, "Revenue Overview", [
        "Q1 revenue: $890K driven by seasonal demand",
        "Q2 revenue: $1.1M boosted by product launch",
        "Q3 revenue: $1.05M with steady enterprise growth",
        "Q4 revenue: $1.16M closing strong with holiday sales",
        "Services revenue grew 22% to represent 30% of total",
    ])

    # Slide 4: Market Trends
    add_content_slide(prs, "Market Trends", [
        "Digital transformation accelerating adoption rates",
        "Competitor pricing pressure in mid-market segment",
        "Emerging markets showing 35% growth potential",
        "AI-powered features becoming key differentiator",
        "Subscription model adoption up 28% year over year",
    ])

    # Slide 5: Regional Performance
    add_content_slide(prs, "Regional Performance", [
        "North America: $2.1M (50% of revenue)",
        "Europe: $1.05M (25% of revenue)",
        "Asia-Pacific: $630K (15% of revenue)",
        "Latin America: $252K (6% of revenue)",
        "Middle East & Africa: $168K (4% of revenue)",
    ])

    # Slide 6: Quarterly Sales Breakdown - EMPTY content area (title only)
    add_title_only_slide(prs, "Quarterly Sales Breakdown")

    # Slide 7: Customer Insights
    add_content_slide(prs, "Customer Insights", [
        "Enterprise clients account for 62% of recurring revenue",
        "Average deal size increased from $45K to $58K",
        "Net Promoter Score improved from 42 to 56",
        "Support ticket resolution time reduced by 35%",
        "Cross-sell conversion rate reached 23%",
    ])

    # Slide 8: Growth Projections
    add_content_slide(prs, "Growth Projections", [
        "2026 revenue target: $5.1M (21% growth)",
        "New product line expected to contribute $600K",
        "International expansion to add 3 new markets",
        "Headcount growth of 15% to support scaling",
        "R&D investment increasing to 18% of revenue",
    ])

    # Slide 9: Thank You
    add_title_slide(prs, "Thank You",
                    "Questions & Discussion\nContact: analytics@company.com")

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
