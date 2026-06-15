"""
Reward Script: Insert stacked bar chart and summary table on slide 6
Task ID: impress_gf2_012
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): Chart exists on slide 6 with stacked bar type
  Component 2 (0.30): Chart has correct 3 data series with expected values
  Component 3 (0.10): Chart has legend and Q1-Q4 categories
  Component 4 (0.15): Table exists on slide 6 with 2 rows x 5 columns
  Component 5 (0.25): Table contains correct header row and totals row
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_012'

# Expected chart data
EXPECTED_SERIES = {
    'ProductA': [30.0, 45.0, 38.0, 52.0],
    'ProductB': [20.0, 25.0, 30.0, 28.0],
    'ProductC': [15.0, 18.0, 22.0, 25.0],
}
EXPECTED_CATEGORIES = ['Q1', 'Q2', 'Q3', 'Q4']
EXPECTED_TABLE_ROW0 = ['Quarter', 'Q1', 'Q2', 'Q3', 'Q4']
EXPECTED_TABLE_ROW1_TOTALS = ['65', '88', '90', '105']  # Q1=65, Q2=88, Q3=90, Q4=105


def persist_app_state(domain):
    """Save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_series_name(series):
    """Extract series name from XML."""
    ns = {'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'}
    tx = series._element.find('.//c:tx', ns)
    if tx is not None:
        v = tx.find('.//c:v', ns)
        if v is not None:
            return v.text
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 6 slides
    if len(prs.slides) < 6:
        print(f"CRITICAL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[5]  # Slide 6 (0-indexed)

    # Find chart and table shapes on slide 6
    chart_shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.CHART]
    table_shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE]

    # Component 1: Chart exists on slide 6 with stacked bar type (0.20 points)
    try:
        if len(chart_shapes) > 0:
            chart = chart_shapes[0].chart
            # BAR_STACKED = 58, also accept COLUMN_STACKED = 52 (both are stacked bar variants)
            chart_type_val = chart.chart_type
            # python-pptx chart_type is an enum; BAR_STACKED=58
            if chart_type_val in (58, 52):  # BAR_STACKED or COLUMN_STACKED
                print(f"PASS: Component 1 — Stacked bar chart found on slide 6 (type={chart_type_val}) (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 1 — Chart found but type is {chart_type_val}, expected stacked bar (58 or 52)")
        else:
            print("FAIL: Component 1 — No chart found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Chart has correct 3 data series with expected values (0.30 points)
    try:
        if len(chart_shapes) > 0:
            chart = chart_shapes[0].chart
            series_list = list(chart.series)
            if len(series_list) == 3:
                series_correct = 0
                for series in series_list:
                    name = get_series_name(series)
                    values = list(series.values)
                    if name in EXPECTED_SERIES:
                        expected_vals = EXPECTED_SERIES[name]
                        if len(values) == len(expected_vals) and all(
                            abs(v - e) < 0.01 for v, e in zip(values, expected_vals)
                        ):
                            series_correct += 1
                            print(f"  Series '{name}' values match: {values}")
                        else:
                            print(f"  Series '{name}' values mismatch: got {values}, expected {expected_vals}")
                    else:
                        print(f"  Series name '{name}' not in expected set")

                if series_correct == 3:
                    print(f"PASS: Component 2 — All 3 series have correct values (0.30 pts)")
                    total_score += 0.30
                elif series_correct >= 1:
                    partial = round(0.30 * series_correct / 3, 2)
                    print(f"PARTIAL: Component 2 — {series_correct}/3 series correct ({partial} pts)")
                    total_score += partial
                else:
                    print(f"FAIL: Component 2 — No series matched expected values")
            else:
                print(f"FAIL: Component 2 — Expected 3 series, found {len(series_list)}")
        else:
            print("FAIL: Component 2 — No chart found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Chart has legend and Q1-Q4 categories (0.10 points)
    try:
        if len(chart_shapes) > 0:
            chart = chart_shapes[0].chart
            has_legend = chart.has_legend
            plot = chart.plots[0]
            categories = [str(c) for c in plot.categories]

            legend_ok = has_legend
            cats_ok = categories == EXPECTED_CATEGORIES

            if legend_ok and cats_ok:
                print(f"PASS: Component 3 — Legend present and categories {categories} correct (0.10 pts)")
                total_score += 0.10
            else:
                if not legend_ok:
                    print(f"FAIL: Component 3 — Chart has no legend")
                if not cats_ok:
                    print(f"FAIL: Component 3 — Categories {categories} != expected {EXPECTED_CATEGORIES}")
        else:
            print("FAIL: Component 3 — No chart found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Table exists on slide 6 with 2 rows x 5 columns (0.15 points)
    try:
        if len(table_shapes) > 0:
            table = table_shapes[0].table
            nrows = len(table.rows)
            ncols = len(table.columns)
            if nrows == 2 and ncols == 5:
                print(f"PASS: Component 4 — Table found with 2x5 dimensions (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 4 — Table dimensions {nrows}x{ncols}, expected 2x5")
        else:
            print("FAIL: Component 4 — No table found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Table contains correct header row and totals row (0.25 points)
    try:
        if len(table_shapes) > 0:
            table = table_shapes[0].table
            nrows = len(table.rows)
            ncols = len(table.columns)

            if nrows >= 2 and ncols >= 5:
                # Check header row
                header = [table.cell(0, c).text.strip() for c in range(5)]
                # Check totals row
                totals = [table.cell(1, c).text.strip() for c in range(5)]

                # Header check: first cell should be something like "Quarter" or header label
                # Data cells should be Q1, Q2, Q3, Q4
                header_q_ok = (header[1:5] == ['Q1', 'Q2', 'Q3', 'Q4'])

                # Totals check: cells 1-4 should match expected totals
                totals_values_ok = (totals[1:5] == EXPECTED_TABLE_ROW1_TOTALS)

                sub_score = 0.0
                if header_q_ok:
                    sub_score += 0.10
                    print(f"  Header quarters match: {header[1:5]}")
                else:
                    print(f"  Header quarters mismatch: got {header[1:5]}, expected ['Q1','Q2','Q3','Q4']")

                if totals_values_ok:
                    sub_score += 0.15
                    print(f"  Totals match: {totals[1:5]}")
                else:
                    print(f"  Totals mismatch: got {totals[1:5]}, expected {EXPECTED_TABLE_ROW1_TOTALS}")

                if sub_score > 0:
                    print(f"PASS: Component 5 — Table content verified ({sub_score} pts)")
                    total_score += sub_score
                else:
                    print(f"FAIL: Component 5 — Table content does not match expected values")
            else:
                print(f"FAIL: Component 5 — Table too small to verify content ({nrows}x{ncols})")
        else:
            print("FAIL: Component 5 — No table found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
