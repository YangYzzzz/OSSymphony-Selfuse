"""
Initial Setup: Create Annual_Review_2025.pptx with 12 slides (no appendix section)
Task ID: impress_ps_048
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_048'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_shape(slide, left, top, width, height, text, font_size=18,
                   bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a text box with formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Annual Review 2025"
    slide1.placeholders[1].text = "Meridian Technologies Inc.\nStrategic Performance Overview"

    # ---- Slide 2: Executive Summary ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_text_shape(slide2, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                   "Executive Summary", font_size=32, bold=True,
                   color=RGBColor(0x1A, 0x23, 0x7E))
    add_text_shape(slide2, Inches(0.5), Inches(1.3), Inches(9), Inches(4),
                   "Meridian Technologies achieved record growth in 2025, "
                   "surpassing revenue targets by 18% and expanding into three new markets. "
                   "Our workforce grew to 2,847 employees across 12 global offices. "
                   "Key milestones included the launch of CloudSync Pro, securing $45M in Series D funding, "
                   "and achieving ISO 27001 certification. Customer satisfaction scores reached an all-time "
                   "high of 94.2%, reflecting our commitment to excellence.",
                   font_size=16)

    # ---- Slide 3: Revenue Overview ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_shape(slide3, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                   "Revenue Overview", font_size=32, bold=True,
                   color=RGBColor(0x1A, 0x23, 0x7E))
    add_text_shape(slide3, Inches(0.5), Inches(1.5), Inches(9), Inches(3.5),
                   "Total Revenue: $128.4M (up 23% YoY)\n"
                   "Recurring Revenue: $89.7M (ARR growth: 31%)\n"
                   "Enterprise Contracts: 247 active accounts\n"
                   "Average Deal Size: $520K (up from $410K)\n"
                   "Net Revenue Retention: 118%",
                   font_size=18)

    # ---- Slide 4: Market Expansion ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_shape(slide4, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                   "Market Expansion", font_size=32, bold=True,
                   color=RGBColor(0x1A, 0x23, 0x7E))
    add_text_shape(slide4, Inches(0.5), Inches(1.5), Inches(9), Inches(4),
                   "New Markets Entered:\n"
                   "  - Southeast Asia (Singapore hub, Q1 2025)\n"
                   "  - Nordic Region (Stockholm office, Q2 2025)\n"
                   "  - Middle East (Dubai partnership, Q3 2025)\n\n"
                   "International Revenue: $41.2M (32% of total)\n"
                   "Cross-border Deals: 78 new enterprise agreements",
                   font_size=16)

    # ---- Slide 5: Product Development ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_shape(slide5, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                   "Product Development", font_size=32, bold=True,
                   color=RGBColor(0x1A, 0x23, 0x7E))
    add_text_shape(slide5, Inches(0.5), Inches(1.5), Inches(9), Inches(4),
                   "Major Releases:\n"
                   "  - CloudSync Pro v3.0 (AI-powered analytics)\n"
                   "  - DataVault Enterprise (zero-trust architecture)\n"
                   "  - MeridianFlow (workflow automation platform)\n\n"
                   "R&D Investment: $19.6M (15.3% of revenue)\n"
                   "Patent Applications Filed: 14\n"
                   "Engineering Headcount: 892 (up 28%)",
                   font_size=16)

    # ---- Slide 6: Customer Success ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_shape(slide6, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                   "Customer Success", font_size=32, bold=True,
                   color=RGBColor(0x1A, 0x23, 0x7E))
    add_text_shape(slide6, Inches(0.5), Inches(1.5), Inches(9), Inches(3.5),
                   "NPS Score: 72 (industry avg: 41)\n"
                   "CSAT: 94.2%\n"
                   "First Response Time: 1.4 hours (down from 3.2)\n"
                   "Resolution Rate: 97.8%\n"
                   "Enterprise Churn: 2.1% (down from 4.7%)\n"
                   "Support Tickets Resolved: 48,320",
                   font_size=18)

    # ---- Slide 7: Team & Culture ----
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_shape(slide7, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                   "Team & Culture", font_size=32, bold=True,
                   color=RGBColor(0x1A, 0x23, 0x7E))
    add_text_shape(slide7, Inches(0.5), Inches(1.5), Inches(9), Inches(4),
                   "Total Employees: 2,847\n"
                   "New Hires in 2025: 634\n"
                   "Employee Engagement Score: 4.6/5.0\n"
                   "Voluntary Turnover: 8.2%\n"
                   "Diversity Index: 0.78\n"
                   "Internal Promotion Rate: 34%\n"
                   "Training Hours per Employee: 42",
                   font_size=18)

    # ---- Slide 8: Financial Highlights ----
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_shape(slide8, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                   "Financial Highlights", font_size=32, bold=True,
                   color=RGBColor(0x1A, 0x23, 0x7E))
    add_text_shape(slide8, Inches(0.5), Inches(1.5), Inches(9), Inches(3.5),
                   "Gross Margin: 72.4%\n"
                   "EBITDA: $31.2M (24.3% margin)\n"
                   "Free Cash Flow: $22.8M\n"
                   "Operating Expenses: $97.2M\n"
                   "Cash & Equivalents: $68.5M\n"
                   "Debt-to-Equity Ratio: 0.32",
                   font_size=18)

    # ---- Slide 9: Strategic Partnerships ----
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_shape(slide9, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                   "Strategic Partnerships", font_size=32, bold=True,
                   color=RGBColor(0x1A, 0x23, 0x7E))
    add_text_shape(slide9, Inches(0.5), Inches(1.5), Inches(9), Inches(4),
                   "Key Partnerships Established:\n"
                   "  - Amazon Web Services (Premier Partner)\n"
                   "  - Deloitte (Systems Integration)\n"
                   "  - Salesforce (AppExchange Launch)\n"
                   "  - SAP (Certified Integration)\n\n"
                   "Partner-Sourced Revenue: $18.4M\n"
                   "Joint Solutions Delivered: 23",
                   font_size=16)

    # ---- Slide 10: Security & Compliance ----
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_shape(slide10, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                   "Security & Compliance", font_size=32, bold=True,
                   color=RGBColor(0x1A, 0x23, 0x7E))
    add_text_shape(slide10, Inches(0.5), Inches(1.5), Inches(9), Inches(4),
                   "Certifications Achieved:\n"
                   "  - ISO 27001 (Information Security)\n"
                   "  - SOC 2 Type II (renewed)\n"
                   "  - GDPR Compliance (verified)\n"
                   "  - HIPAA Ready (healthcare vertical)\n\n"
                   "Security Incidents: 0 critical breaches\n"
                   "Penetration Tests Passed: 4\n"
                   "Uptime SLA: 99.97%",
                   font_size=16)

    # ---- Slide 11: 2026 Roadmap ----
    slide11 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_shape(slide11, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                   "2026 Roadmap", font_size=32, bold=True,
                   color=RGBColor(0x1A, 0x23, 0x7E))
    add_text_shape(slide11, Inches(0.5), Inches(1.5), Inches(9), Inches(4),
                   "Strategic Priorities:\n"
                   "  - Launch AI Assistant Suite (Q1)\n"
                   "  - Expand APAC presence to 5 countries\n"
                   "  - Achieve $180M revenue target\n"
                   "  - IPO readiness assessment (Q3)\n"
                   "  - Hire 400+ engineering talent\n"
                   "  - Launch Meridian University training program",
                   font_size=16)

    # ---- Slide 12: Closing / Thank You ----
    slide12 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = slide12.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    add_text_shape(slide12, Inches(1), Inches(2.5), Inches(8), Inches(1.5),
                   "Thank You", font_size=44, bold=True,
                   color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_text_shape(slide12, Inches(1), Inches(4.0), Inches(8), Inches(1),
                   "Questions & Discussion", font_size=24,
                   color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
