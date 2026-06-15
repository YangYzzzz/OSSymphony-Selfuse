"""
Reward Script: Create appendix section in Annual_Review_2025.pptx
Task ID: impress_ps_048
Domain: libreoffice_impress
Scoring:
  C1: Slide count == 14                                         (0.10)
  C2: Slide 12 has dark blue (#1A237E) background               (0.15)
  C3: Slide 12 has 'Appendix' centered, bold, white text        (0.20)
  C4: Slide 13 has table 13r x 5c                               (0.20)
  C5: Table headers correct with bold white formatting           (0.15)
  C6: Table has 12 month rows (Jan-Dec) with numeric data       (0.10)
  C7: Slide 14 is the original closing slide ('Thank You')      (0.10)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_048'

EXPECTED_MONTHS = [
    'January', 'February', 'March', 'April', 'May', 'June',
    'July', 'August', 'September', 'October', 'November', 'December'
]

EXPECTED_HEADERS = ['Month', 'Revenue', 'Expenses', 'Net', 'Growth%']


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Presentation has exactly 14 slides (0.10 pts)
    # Initial has 12 slides; golden should have 14 (added divider + table slide)
    try:
        if num_slides == 14:
            print(f"PASS: Component 1 -- slide count is 14 (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 1 -- expected 14 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Gate: need at least 13 slides to check the new ones
    if num_slides < 13:
        print(f"CRITICAL: Not enough slides ({num_slides}) to verify appendix content")
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Slide 12 has dark blue (#1A237E) solid background (0.15 pts)
    # Initial slide 12 is "Thank You" which already has 1A237E bg, but in initial
    # there are only 12 slides and slide 12 is the closing. In golden, slide 12
    # is a NEW divider slide. Since we gated on num_slides >= 13, this only fires
    # when new slides were added.
    try:
        slide12 = prs.slides[11]  # 0-indexed
        fill = slide12.background.fill
        if fill.type == 1:  # SOLID
            bg_color = str(fill.fore_color.rgb).upper()
            if bg_color == '1A237E':
                print(f"PASS: Component 2 -- slide 12 bg is #1A237E (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 2 -- slide 12 bg is #{bg_color}, expected #1A237E")
        else:
            print(f"FAIL: Component 2 -- slide 12 bg fill type is {fill.type}, expected SOLID (1)")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Slide 12 has 'Appendix' text, bold, white, centered (0.20 pts)
    # This is the key differentiator -- initial has no 'Appendix' text on any slide
    try:
        slide12 = prs.slides[11]
        appendix_found = False
        appendix_bold = False
        appendix_white = False
        appendix_centered = False

        for shape in slide12.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    full_text = para.text.strip().lower()
                    if 'appendix' in full_text:
                        appendix_found = True
                        # Check alignment: CENTER (2) expected
                        if para.alignment == PP_ALIGN.CENTER:
                            appendix_centered = True
                        for run in para.runs:
                            if 'appendix' in run.text.strip().lower():
                                if run.font.bold is True:
                                    appendix_bold = True
                                try:
                                    if str(run.font.color.rgb).upper() == 'FFFFFF':
                                        appendix_white = True
                                except:
                                    pass

        sub_score = 0.0
        if appendix_found:
            sub_score += 0.08
        if appendix_bold:
            sub_score += 0.04
        if appendix_white:
            sub_score += 0.04
        if appendix_centered:
            sub_score += 0.04

        if sub_score > 0:
            total_score += sub_score
            details = f"found={appendix_found}, bold={appendix_bold}, white={appendix_white}, centered={appendix_centered}"
            print(f"PASS: Component 3 -- Appendix text on slide 12 ({sub_score:.2f} pts) [{details}]")
        else:
            print(f"FAIL: Component 3 -- No 'Appendix' text found on slide 12")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Slide 13 has a table with 13 rows x 5 columns (0.20 pts)
    # Initial has no table on any slide, so this is purely task-introduced
    try:
        slide13 = prs.slides[12]
        table_found = False
        correct_dims = False

        for shape in slide13.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                table_found = True
                rows = len(table.rows)
                cols = len(table.columns)
                if rows == 13 and cols == 5:
                    correct_dims = True
                    print(f"PASS: Component 4 -- table 13x5 found on slide 13 (0.20 pts)")
                    total_score += 0.20
                else:
                    # Partial: table exists but wrong dimensions
                    total_score += 0.05
                    print(f"PARTIAL: Component 4 -- table found but {rows}x{cols}, expected 13x5 (0.05 pts)")
                break

        if not table_found:
            print(f"FAIL: Component 4 -- no table found on slide 13")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Table headers match expected with bold white formatting (0.15 pts)
    # Headers: Month, Revenue, Expenses, Net, Growth%
    try:
        slide13 = prs.slides[12]
        table_obj = None
        for shape in slide13.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_obj = shape.table
                break

        if table_obj is not None and len(table_obj.columns) >= 5:
            headers_match = 0
            header_formatting = 0

            for c in range(min(5, len(table_obj.columns))):
                cell_text = table_obj.cell(0, c).text.strip()
                if cell_text == EXPECTED_HEADERS[c]:
                    headers_match += 1

                # Check bold + white on header runs
                for para in table_obj.cell(0, c).text_frame.paragraphs:
                    for run in para.runs:
                        is_bold = run.font.bold is True
                        is_white = False
                        try:
                            is_white = str(run.font.color.rgb).upper() == 'FFFFFF'
                        except:
                            pass
                        if is_bold and is_white:
                            header_formatting += 1
                        break
                    break

            sub_score = 0.0
            if headers_match == 5:
                sub_score += 0.09
            elif headers_match >= 3:
                sub_score += 0.04

            if header_formatting >= 4:
                sub_score += 0.06
            elif header_formatting >= 2:
                sub_score += 0.03

            if sub_score > 0:
                total_score += sub_score
                print(f"PASS: Component 5 -- headers: {headers_match}/5 match, {header_formatting}/5 formatted ({sub_score:.2f} pts)")
            else:
                print(f"FAIL: Component 5 -- headers: {headers_match}/5 match, {header_formatting}/5 formatted")
        else:
            print(f"FAIL: Component 5 -- no suitable table found on slide 13")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # Component 6: Table has 12 month rows (Jan-Dec) with numeric data (0.10 pts)
    try:
        slide13 = prs.slides[12]
        table_obj = None
        for shape in slide13.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_obj = shape.table
                break

        if table_obj is not None and len(table_obj.rows) >= 13:
            months_found = 0
            data_filled = 0

            for r in range(1, min(13, len(table_obj.rows))):
                month_text = table_obj.cell(r, 0).text.strip()
                if month_text in EXPECTED_MONTHS:
                    months_found += 1

                # Check that other cells have some data (non-empty)
                row_has_data = all(
                    table_obj.cell(r, c).text.strip() != ''
                    for c in range(1, min(5, len(table_obj.columns)))
                )
                if row_has_data:
                    data_filled += 1

            sub_score = 0.0
            if months_found == 12:
                sub_score += 0.05
            elif months_found >= 6:
                sub_score += 0.02

            if data_filled >= 10:
                sub_score += 0.05
            elif data_filled >= 6:
                sub_score += 0.02

            if sub_score > 0:
                total_score += sub_score
                print(f"PASS: Component 6 -- {months_found}/12 months, {data_filled}/12 data rows ({sub_score:.2f} pts)")
            else:
                print(f"FAIL: Component 6 -- {months_found}/12 months, {data_filled}/12 data rows")
        else:
            print(f"FAIL: Component 6 -- table not found or too few rows")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    # Component 7: Slide 14 (last slide) is the original closing slide with 'Thank You' (0.10 pts)
    # In initial, this was slide 12. In golden, it should be slide 14.
    try:
        if num_slides >= 14:
            slide14 = prs.slides[13]
            thank_you_found = False
            for shape in slide14.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip().lower()
                    if 'thank you' in text:
                        thank_you_found = True
                        break

            if thank_you_found:
                print(f"PASS: Component 7 -- slide 14 has 'Thank You' (original closing) (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 7 -- slide 14 does not contain 'Thank You'")
        else:
            print(f"FAIL: Component 7 -- only {num_slides} slides, expected at least 14")
    except Exception as e:
        print(f"ERROR: Component 7 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook: save any unsaved LibreOffice state before verification
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
