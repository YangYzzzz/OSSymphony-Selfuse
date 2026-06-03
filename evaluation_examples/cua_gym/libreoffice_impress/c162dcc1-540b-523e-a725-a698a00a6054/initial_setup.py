"""
Initial Setup: Single slide ad presentation with no transitions
Task ID: impress_tm_024
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_024'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Special Offer ad ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Background - dark blue/navy for an ad look
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)

    # Title: "Special Offer"
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(0.8), Inches(10), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Special Offer"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xD7, 0x00)  # Gold

    # Subtitle / tagline
    txBox2 = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Limited Time Only — Up to 40% Off Select Items"
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.name = "Arial"
    run2.font.size = Pt(28)
    run2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Promo details
    txBox3 = slide.shapes.add_textbox(Inches(2), Inches(3.8), Inches(9), Inches(2.5))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True

    details = [
        "Premium Wireless Headphones — $89.99 (was $149.99)",
        "Smart Fitness Tracker Pro — $59.99 (was $99.99)",
        "Ultra-Slim Power Bank 20000mAh — $29.99 (was $49.99)",
        "Bluetooth Speaker System — $119.99 (was $199.99)",
    ]

    for i, detail in enumerate(details):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = detail
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    # Call to action
    txBox4 = slide.shapes.add_textbox(Inches(3), Inches(6.0), Inches(7), Inches(0.8))
    tf4 = txBox4.text_frame
    tf4.word_wrap = True
    p4 = tf4.paragraphs[0]
    p4.text = "Visit store.example.com | Use code SAVE40"
    p4.alignment = PP_ALIGN.CENTER
    run4 = p4.runs[0]
    run4.font.name = "Arial"
    run4.font.size = Pt(20)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0xFF, 0xD7, 0x00)

    # No transitions, no auto-advance — default is advance on click
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
