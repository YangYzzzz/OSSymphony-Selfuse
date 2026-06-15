"""
FINAL REWARD SCRIPT - SUCCESS
Task: In LibreOffice Impress, what’s the quickest way to turn only slides 1 through 5 into a PDF and drop the file on my Desktop as “intro-1-5.pdf”?
Generated: 2025-09-11 00:24:23
Status: success
Model: azure-o3
Total Steps: 18
"""

import os
import re
from pptx import Presentation
import PyPDF2


def _extract_first_text_lines(presentation, slide_numbers):
    """Extract the first meaningful text line from each requested slide.

    Parameters
    ----------
    presentation : pptx.Presentation
        Loaded PPTX presentation.
    slide_numbers : list[int]
        1-based slide indices to inspect.

    Returns
    -------
    list[str | None]
        List of extracted text fragments (None when no text found on a slide).
    """
    texts = []
    for idx in slide_numbers:
        if 1 <= idx <= len(presentation.slides):
            slide = presentation.slides[idx - 1]
            first_line = None
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    first_line = re.sub(r"\s+", " ", shape.text.strip())
                    break
            texts.append(first_line)
        else:
            texts.append(None)
    return texts


def verify_intro_pdf():
    """Verify that slides 1-5 have been exported to a PDF on the Desktop.

    Returns
    -------
    float
        Reward score between 0.0 and 1.0
    """
    score = 0.0
    max_score = 1.0

    # Paths
    pdf_path = os.path.expanduser("~/Desktop/intro-1-5.pdf")
    pptx_path = (
        "/home/user/"
        "in_libreoffice_impress_whats_the_quickest_way_to_turn_only_slides_1_"
        "through_5_into_a_pdf_and_drop_th_golden.pptx"
    )

    # ------------------------------------------------------------------
    # 1) Verify the exported PDF file exists on the Desktop
    # ------------------------------------------------------------------
    if not os.path.exists(pdf_path):
        print(f"✗ Expected PDF not found at {pdf_path}")
        print("REWARD:", score)
        return score  # 0.0 – nothing else to check
    print(f"✓ PDF found: {pdf_path}")

    # ------------------------------------------------------------------
    # 2) Verify the PDF has exactly 5 pages (slides 1-5)
    # ------------------------------------------------------------------
    try:
        reader = PyPDF2.PdfReader(pdf_path)
        page_count = len(reader.pages)
        print(f"PDF page count: {page_count}")

        if page_count == 5:
            score += 0.5  # half the points for correct page range
            print("✓ PDF has correct number of pages (0.5 points)")
        else:
            print("✗ PDF does not have exactly 5 pages; 0.0 points for page count")
    except Exception as e:
        print(f"✗ Could not read PDF: {e}")
        print("REWARD:", score)
        return score  # Cannot continue without reading the PDF

    # ------------------------------------------------------------------
    # 3) Load source PPTX and gather expected text for slides 1-5
    # ------------------------------------------------------------------
    try:
        prs = Presentation(pptx_path)
        expected_texts = _extract_first_text_lines(prs, [1, 2, 3, 4, 5])
    except Exception as e:
        print(f"✗ Error loading PPTX or extracting texts: {e}")
        expected_texts = [None] * 5  # fall back to None so no matches awarded

    # ------------------------------------------------------------------
    # 4) Compare text content of each PDF page with the corresponding slide
    # ------------------------------------------------------------------
    matched_pages = 0
    pages_to_check = min(5, len(reader.pages))

    for idx in range(pages_to_check):
        try:
            page = reader.pages[idx]
            extracted_text = page.extract_text() or ""
            extracted_norm = re.sub(r"\s+", " ", extracted_text.strip())
            expected_fragment = expected_texts[idx]

            if expected_fragment and expected_fragment in extracted_norm:
                matched_pages += 1
                print(f"✓ Page {idx + 1} matches expected text '{expected_fragment}'")
            else:
                print(
                    f"✗ Page {idx + 1} text mismatch. "
                    f"Expected fragment: '{expected_fragment}', "
                    f"Extracted: '{extracted_norm[:60]}'"
                )
        except Exception as e:
            print(f"✗ Error extracting text from page {idx + 1}: {e}")

    # Award remaining 0.5 points proportionally for text matches
    text_match_score = 0.5 * (matched_pages / 5)
    score += text_match_score

    if matched_pages == 5:
        print("✓ All 5 pages match the expected slide content (0.5 points)")
    else:
        print(
            f"Partial match: {matched_pages}/5 pages matched "
            f"({text_match_score:.2f} points)"
        )

    # ------------------------------------------------------------------
    # 5) Finalise & cap score
    # ------------------------------------------------------------------
    final_score = round(min(score, max_score), 2)
    print(f"Final Score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score


if __name__ == "__main__":
    verify_intro_pdf()

