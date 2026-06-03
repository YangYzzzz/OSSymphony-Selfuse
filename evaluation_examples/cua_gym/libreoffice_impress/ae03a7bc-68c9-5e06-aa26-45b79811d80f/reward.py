"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 92, I need to pop in a perfect circle that’s exactly 4 cm tall and 4 cm wide, filled with the “Purple 2” swatch from the Impress palette. What’s the quickest way to do that?
Generated: 2025-09-10 17:40:46
Status: success
Model: azure-o3
Total Steps: 8
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE

# -------------------------------------------------------------
# Helper: Convert centimetres to English Metric Units (EMU)
# -------------------------------------------------------------
EMU_PER_CM = 360000  # constant used by PowerPoint

def emu_from_cm(cm: float) -> int:
    """Return EMU value for the given centimetres (rounded)."""
    return int(round(cm * EMU_PER_CM))

# -------------------------------------------------------------
# Main verification routine
# -------------------------------------------------------------

def verify_circle_on_slide_92(file_path: str) -> float:
    """Verify slide-92 contains a 4 cm × 4 cm purple circle.

    Scoring (progressive):
      • 0.4 – at least one circle (oval auto-shape) present on slide 92
      • 0.3 – at least one circle is exactly 4 cm × 4 cm (±0.02 cm tolerance)
      • 0.3 – that circle is filled with the required “Purple 2” colour

    A perfect submission earns 1.0; partial completions earn partial credit.
    """
    total_score = 0.0
    print(f"Verifying file: {file_path}\n")

    # ---------- prerequisite checks ----------
    if not os.path.exists(file_path):
        print("✗ File not found – cannot verify task.")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Could not open presentation: {e}")
        return 0.0

    # ---------- requirement 1: slide 92 must exist ----------
    slide_index = 91  # zero-based index
    if len(prs.slides) <= slide_index:
        print(f"✗ Presentation has only {len(prs.slides)} slides – slide 92 missing.")
        return 0.0
    print(f"✓ Presentation contains {len(prs.slides)} slides.")

    slide92 = prs.slides[slide_index]

    # ---------- requirement 2: an oval (circle) must be present ----------
    circles = [sh for sh in slide92.shapes
               if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and
                  sh.auto_shape_type == MSO_AUTO_SHAPE_TYPE.OVAL]

    if not circles:
        print("✗ No circle (oval auto-shape) found on slide 92.")
        return 0.0

    print(f"✓ Found {len(circles)} oval auto-shape(s) on slide 92 (0.4 points).")
    total_score += 0.4

    # ---------- requirement 3: check size (4 cm × 4 cm) ----------
    TARGET_EMU = emu_from_cm(4.0)
    TOLERANCE_EMU = 20000  # ≈0.02 cm

    size_ok = False
    for sh in circles:
        w_ok = abs(sh.width  - TARGET_EMU) <= TOLERANCE_EMU
        h_ok = abs(sh.height - TARGET_EMU) <= TOLERANCE_EMU
        aspect_ok = abs(sh.width - sh.height) <= TOLERANCE_EMU
        if w_ok and h_ok and aspect_ok:
            size_ok = True
            print(f"  ✓ Circle with correct size detected (w={sh.width}, h={sh.height}).")
            break
        else:
            print(f"  – Circle size mismatch (w={sh.width}, h={sh.height}).")

    if size_ok:
        total_score += 0.3
    else:
        print("✗ No circle matched the exact 4 cm × 4 cm requirement.")

    # ---------- requirement 4: check fill colour (Purple 2) ----------
    # Known RGB hex values for Impress palette “Purple 2”.
    PURPLE2_RGB_HEX = {
        "9A32CD",  # observed in golden file
        "9A33CD"   # possible rounding variation
    }

    color_ok = False
    for sh in circles:
        # Only inspect solid fills
        if sh.fill.type == 1:  # SOLID
            rgb_object = sh.fill.fore_color.rgb
            if rgb_object is None:
                continue
            rgb_hex = str(rgb_object).upper()
            print(f"  Detected fill RGB: {rgb_hex}")
            if rgb_hex in PURPLE2_RGB_HEX:
                color_ok = True
                break

    if color_ok:
        print("✓ Circle fill colour matches Purple 2 (0.3 points).")
        total_score += 0.3
    else:
        print("✗ Circle fill colour does not match Purple 2.")

    # ---------- final score ----------
    final_score = min(total_score, 1.0)
    print(f"\nTotal score: {final_score}")
    return final_score

# -------------------------------------------------------------
# Run verification when script is executed
# -------------------------------------------------------------
if __name__ == "__main__":
    test_path = "/home/user/on_slide_92_i_need_to_pop_in_a_perfect_circle_thats_exactly_4_cm_tall_and_4_cm_wide_filled_with_the__golden.pptx"
    reward_value = verify_circle_on_slide_92(test_path)
    print(f"REWARD: {reward_value}")
