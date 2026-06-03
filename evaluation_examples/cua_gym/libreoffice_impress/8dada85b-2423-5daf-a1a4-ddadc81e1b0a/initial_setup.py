"""
Initial Setup: Leadership Talk presentation with 8 slides, no notes on slides 3-5,
presenter console disabled.
Task ID: impress_gf2_048
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_048'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_body(slide, title_text, body_lines):
    """Helper to populate title and body placeholders."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    # Find body placeholder (idx 1 typically)
    body_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            body_ph = ph
            break
    if body_ph and body_ph.has_text_frame:
        tf = body_ph.text_frame
        tf.clear()
        for i, line in enumerate(body_lines):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                p = tf.add_paragraph()
                p.text = line


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "The Art of Leadership"
    slide1.placeholders[1].text = "Strategies for Inspiring Teams and Driving Results\nPresented by Dr. Amara Okafor\nGlobal Leadership Summit 2025"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide2, "Session Agenda", [
        "1. Defining Modern Leadership",
        "2. Building Trust Within Your Team",
        "3. Communication Strategies That Work",
        "4. Managing Change and Uncertainty",
        "5. Emotional Intelligence in the Workplace",
        "6. Case Studies and Real-World Examples",
        "7. Q&A and Action Planning",
    ])

    # --- Slide 3: Building Trust (NO NOTES) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide3, "Building Trust Within Your Team", [
        "Trust is the foundation of effective leadership",
        "Consistency between words and actions",
        "Transparent decision-making processes",
        "Acknowledging mistakes openly",
        "Empowering team members with autonomy",
        "Creating psychological safety for honest feedback",
    ])

    # --- Slide 4: Communication Strategies (NO NOTES) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide4, "Communication Strategies That Work", [
        "Active listening: hear to understand, not to respond",
        "Tailoring messages to your audience",
        "Using storytelling to convey complex ideas",
        "Non-verbal cues and body language awareness",
        "Regular one-on-one check-ins with direct reports",
        "Feedback loops: giving and receiving constructively",
    ])

    # --- Slide 5: Managing Change (NO NOTES) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide5, "Managing Change and Uncertainty", [
        "The ADKAR model for change management",
        "Communicating the 'why' behind organizational shifts",
        "Supporting employees through transition periods",
        "Building resilience within teams",
        "Measuring adoption and addressing resistance",
        "Celebrating milestones during transformation",
    ])

    # --- Slide 6: Emotional Intelligence ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide6, "Emotional Intelligence in Leadership", [
        "Self-awareness: understanding your triggers",
        "Self-regulation: responding vs. reacting",
        "Motivation: connecting work to purpose",
        "Empathy: seeing from others' perspectives",
        "Social skills: navigating complex team dynamics",
    ])
    slide6.notes_slide.notes_text_frame.text = "Mention Daniel Goleman's framework briefly."

    # --- Slide 7: Case Studies ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide7, "Case Studies: Leadership in Action", [
        "Case 1: Turnaround at Meridian Healthcare",
        "  - CEO rebuilt trust after compliance scandal",
        "  - Transparent town halls, open-door policy",
        "Case 2: Rapid Growth at NovaTech Solutions",
        "  - VP Engineering scaled team from 12 to 85",
        "  - Maintained culture through structured mentorship",
    ])
    slide7.notes_slide.notes_text_frame.text = "Ask audience which case resonates more with their experience."

    # --- Slide 8: Q&A and Closing ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide8, "Q&A and Action Planning", [
        "What is one leadership practice you will adopt this week?",
        "Write down your commitment and share with a partner",
        "Resources and recommended reading list available online",
        "Contact: a.okafor@leadershipsummit.org",
        "Thank you for your engagement today!",
    ])

    # Do NOT add notes to slides 3, 4, 5 (task requires agent to add them)
    # Do NOT enable presenter console (task requires agent to enable it)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
