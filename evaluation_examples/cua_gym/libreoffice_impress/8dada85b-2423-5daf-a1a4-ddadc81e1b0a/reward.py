"""
Reward Script: Configure speaker notes for slides 3, 4, 5 and enable presenter console
Task ID: impress_gf2_048
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Slide 3 has speaker notes with >= 3 sentences
  Component 2 (0.25): Slide 4 has speaker notes with >= 3 sentences
  Component 3 (0.25): Slide 5 has speaker notes with >= 3 sentences
  Component 4 (0.25): showPr element present in pptx XML (presenter view config)
"""

import os
import re
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_048'


def count_sentences(text):
    """Count sentences in text. A sentence ends with '.', '!', or '?'."""
    if not text or not text.strip():
        return 0
    # Split on sentence-ending punctuation followed by space or end of string
    sentences = re.split(r'[.!?]+(?:\s|$)', text.strip())
    # Filter out empty strings from the split
    sentences = [s.strip() for s in sentences if s.strip()]
    return len(sentences)


def get_slide_notes(prs, slide_idx):
    """Get notes text for a slide (0-indexed). Returns empty string if no notes."""
    try:
        slide = prs.slides[slide_idx]
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def check_show_pr(file_path):
    """Check if showPr element exists in presProps.xml or presentation.xml."""
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            # Check presProps.xml
            for xml_name in ['ppt/presProps.xml', 'ppt/presentation.xml']:
                if xml_name in zf.namelist():
                    with zf.open(xml_name) as f:
                        content = f.read().decode()
                        # Look for showPr element in the XML
                        if 'showPr' in content:
                            return True
        return False
    except Exception:
        return False


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"CRITICAL: Expected at least 5 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide 3 has speaker notes with >= 3 sentences (0.25 points)
    try:
        notes_3 = get_slide_notes(prs, 2)  # 0-indexed
        sentence_count_3 = count_sentences(notes_3)
        if len(notes_3) >= 50 and sentence_count_3 >= 3:
            print(f"PASS: Component 1 -- Slide 3 has {sentence_count_3} sentences, {len(notes_3)} chars (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 -- Slide 3 notes: {sentence_count_3} sentences, {len(notes_3)} chars (need >= 3 sentences and >= 50 chars)")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Slide 4 has speaker notes with >= 3 sentences (0.25 points)
    try:
        notes_4 = get_slide_notes(prs, 3)  # 0-indexed
        sentence_count_4 = count_sentences(notes_4)
        if len(notes_4) >= 50 and sentence_count_4 >= 3:
            print(f"PASS: Component 2 -- Slide 4 has {sentence_count_4} sentences, {len(notes_4)} chars (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 -- Slide 4 notes: {sentence_count_4} sentences, {len(notes_4)} chars (need >= 3 sentences and >= 50 chars)")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Slide 5 has speaker notes with >= 3 sentences (0.25 points)
    try:
        notes_5 = get_slide_notes(prs, 4)  # 0-indexed
        sentence_count_5 = count_sentences(notes_5)
        if len(notes_5) >= 50 and sentence_count_5 >= 3:
            print(f"PASS: Component 3 -- Slide 5 has {sentence_count_5} sentences, {len(notes_5)} chars (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 -- Slide 5 notes: {sentence_count_5} sentences, {len(notes_5)} chars (need >= 3 sentences and >= 50 chars)")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: showPr element present in pptx XML - indicates presenter view configured (0.25 points)
    try:
        has_show_pr = check_show_pr(file_path)
        if has_show_pr:
            print(f"PASS: Component 4 -- showPr element found in pptx XML (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 4 -- showPr element not found in pptx XML")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
