"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 20’s heading needs a facelift: swap whatever font it’s wearing now for Liberation Serif at 40 pt and make sure the whole title is underlined. How do I do that in LibreOffice Impress?
Generated: 2025-09-10 18:55:05
Status: success
Model: azure-o3
Total Steps: 5
"""

import os
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import PP_PLACEHOLDER

"""
Reward Script for LibreOffice Impress Task
Task: On slide 20, the title (heading) must
  1. Use the font "Liberation Serif"
  2. Have a font size of 40 pt
  3. Be completely underlined
The script awards progressive points for each of the three independent
requirements (0.4 + 0.3 + 0.3 = 1.0 total).
It verifies every run of the title text to ensure **all** runs satisfy the
conditions – partial compliance does not earn points for that criterion.
"""

FILE_PATH = "/home/user/slide_20s_heading_needs_a_facelift_swap_whatever_font_its_wearing_now_for_liberation_serif_at_40_pt__golden.pptx"


def _locate_title_shape(slide):
    """Return the title placeholder if present, otherwise the first text shape."""
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
            return shape
    for shape in slide.shapes:
        if shape.has_text_frame:
            return shape
    return None


def verify_slide20_heading(file_path):
    print(f"Verifying presentation: {file_path}")

    # ---------- Basic file checks (no points awarded) ----------
    if not os.path.exists(file_path):
        print("✗ File not found")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print("✗ Error loading presentation:", e)
        return 0.0

    if len(prs.slides) < 20:
        print("✗ Presentation contains fewer than 20 slides")
        return 0.0

    # ---------- Locate title shape on slide 20 ----------
    slide20 = prs.slides[19]  # zero-indexed
    title_shape = _locate_title_shape(slide20)
    if title_shape is None:
        print("✗ Could not locate a title/text shape on slide 20")
        return 0.0
    print(f"✓ Found title shape with text: '{title_shape.text.strip()}'")

    # ---------- Examine all runs in the title ----------
    correct_font = True
    correct_size = True
    correct_underline = True
    run_count = 0

    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run_count += 1
            name = (run.font.name or "").strip()
            size = run.font.size  # EMU integer or None
            underline = run.font.underline

            print(f"  Run '{run.text}' -> font='{name}', size={size}, underline={underline}")

            # Font name check
            if name.lower() != "liberation serif":
                correct_font = False

            # Font size check (size is stored as EMUs: 1pt = 12700 EMU)
            target_size = Pt(40)
            if size is None or abs(size - target_size) > Pt(0.5):  # 0.5 pt tolerance
                correct_size = False

            # Underline check (True/False/None)
            if not underline:
                correct_underline = False

    if run_count == 0:
        print("✗ No text runs found in the title – cannot verify")
        return 0.0

    # ---------- Progressive Scoring ----------
    score = 0.0

    if correct_font:
        score += 0.4
        print("✓ All runs use Liberation Serif (0.4)")
    else:
        print("✗ Font name requirement failed (0.0)")

    if correct_size:
        score += 0.3
        print("✓ All runs are 40 pt (0.3)")
    else:
        print("✗ Font size requirement failed (0.0)")

    if correct_underline:
        score += 0.3
        print("✓ All runs are underlined (0.3)")
    else:
        print("✗ Underline requirement failed (0.0)")

    final = min(score, 1.0)
    print(f"Total score: {final}")
    return final


# --------------------- Run Verification ---------------------
reward = verify_slide20_heading(FILE_PATH)
print(f"REWARD: {reward}")
