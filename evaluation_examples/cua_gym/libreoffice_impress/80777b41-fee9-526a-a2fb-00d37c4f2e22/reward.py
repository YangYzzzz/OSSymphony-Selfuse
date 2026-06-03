"""
FINAL REWARD SCRIPT - SUCCESS
Task: I’m tidying up my LibreOffice Impress presentation and want the very first slide to pop. On Slide 1 please make every piece of text bold, and set the title itself to 44 pt with an underline. How do I do that?
Generated: 2025-09-10 14:02:03
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation
from pptx.util import Pt

"""
Reward Script for LibreOffice Impress / PowerPoint Task
Task:  
Make every piece of text on Slide 1 bold, and set the title text to
(1) 44 pt font size and (2) underlined.

Scoring (progressive, 0‒1):
 • Bold for ALL text runs on Slide 1              → 0.60 pts (60 %)
 • Title font size = 44 pt                        → 0.20 pts (20 %)
 • Title text underlined                          → 0.20 pts (20 %)

The script inspects Slide 1 of the given .pptx file and awards points
ONLY when the specified formatting requirements are actually met.
It prints detailed diagnostics and finally prints "REWARD: X.X" where
X.X ∈ [0.0, 1.0].
"""

FILE_PATH = (
    "/home/user/"
    "im_tidying_up_my_libreoffice_impress_presentation_and_want_the_"
    "very_first_slide_to_pop_on_slide_1_pl_golden.pptx"
)

def verify_slide1_text_format(file_path: str) -> float:
    """Verify Slide 1 formatting according to task requirements."""

    print(f"Loading presentation: {file_path}")

    # ---- Basic file checks (no points awarded) --------------------------
    if not os.path.exists(file_path):
        print("✗ File not found")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as exc:
        print(f"✗ Error loading presentation: {exc}")
        return 0.0

    if len(prs.slides) == 0:
        print("✗ Presentation has no slides")
        return 0.0

    slide = prs.slides[0]  # Slide 1 (0-indexed)

    # --------------------------------------------------------------------
    # Requirement 1: Every text run on Slide 1 is bold → 0.60 pts
    # --------------------------------------------------------------------
    total_runs = 0
    bold_runs = 0

    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or shape.text_frame is None:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text is None or run.text.strip() == "":
                    continue  # ignore empty runs
                total_runs += 1
                if run.font.bold is True:
                    bold_runs += 1

    if total_runs == 0:
        print("✗ No text runs found on Slide 1 → cannot award bold points")
        bold_score = 0.0
    else:
        bold_ratio = bold_runs / total_runs
        bold_score = round(bold_ratio * 0.60, 4)  # up to 0.60 pts
        print(
            f"✓ Bold verification: {bold_runs}/{total_runs} runs bold "
            f"→ {bold_score:.2f} pts"
        )

    # --------------------------------------------------------------------
    # Requirements 2 & 3: Title formatting (size & underline)
    # Each worth 0.20 pts (total 0.40)
    # --------------------------------------------------------------------
    title_shape = slide.shapes.title
    size_score = underline_score = 0.0

    if title_shape is None or title_shape.text_frame is None:
        print("✗ No title placeholder on Slide 1 → cannot check title formatting")
    else:
        title_runs = [
            run
            for para in title_shape.text_frame.paragraphs
            for run in para.runs
            if run.text and run.text.strip()
        ]

        if not title_runs:
            print("✗ Title placeholder contains no text → no title points")
        else:
            # Font size check (allow tiny tolerance)
            correct_size_runs = sum(
                1
                for run in title_runs
                if run.font.size is not None and abs(run.font.size - Pt(44)) <= 1_000  # EMU tolerance
            )
            size_score = round(
                (correct_size_runs / len(title_runs)) * 0.20, 4
            )
            if size_score == 0:
                print("✗ Title font size ≠ 44 pt for all runs")
            else:
                print(
                    f"✓ Title size: {correct_size_runs}/{len(title_runs)} runs at 44 pt "
                    f"→ {size_score:.2f} pts"
                )

            # Underline check
            underlined_runs = sum(
                1 for run in title_runs if run.font.underline is True
            )
            underline_score = round(
                (underlined_runs / len(title_runs)) * 0.20, 4
            )
            if underline_score == 0:
                print("✗ Title text not fully underlined")
            else:
                print(
                    f"✓ Title underline: {underlined_runs}/{len(title_runs)} runs underlined "
                    f"→ {underline_score:.2f} pts"
                )

    # ---- Total score ----------------------------------------------------
    total_score = round(min(bold_score + size_score + underline_score, 1.0), 4)
    print(f"Total score: {total_score}")
    return total_score


if __name__ == "__main__":
    reward = verify_slide1_text_format(FILE_PATH)
    print(f"REWARD: {reward}")

