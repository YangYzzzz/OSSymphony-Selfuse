"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 295 is the only one that hasn’t been updated yet. In LibreOffice Impress, how do I switch its background to Green 2 (#00A933) and set the title text to pure white (#FFFFFF)?
Generated: 2025-09-10 19:43:56
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE


def verify_slide_295_background_and_title(file_path: str) -> float:
    """Verify that slide 295 (index 294) has:
    1. Background set to Green 2 (#00A933)
    2. Title text colour set to pure white (#FFFFFF)

    Returns a progressive score between 0.0 and 1.0.
    """

    print(f"Verifying presentation: {file_path}")
    max_score = 1.0
    score = 0.0

    # --- Preliminary checks -------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File not found.")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as exc:
        print(f"✗ Failed to load presentation: {exc}")
        return 0.0

    total_slides = len(prs.slides)
    print(f"Total slides detected: {total_slides}")

    if total_slides < 295:
        print("✗ Expected at least 295 slides (need slide index 294).")
        return 0.0

    target_slide = prs.slides[294]  # Slide 295 (0-based index)

    # --- Requirement 1: Background colour is Green 2 ------------------------
    try:
        bg_fill = target_slide.background.fill
        bg_ok = False

        if bg_fill.type == MSO_FILL_TYPE.SOLID:
            fore_colour = bg_fill.fore_color
            # colour type 1 == RGB
            if getattr(fore_colour, "type", None) == 1:  # RGB colour present
                rgb = str(fore_colour.rgb).upper()
                print(f"Background RGB detected: {rgb}")
                if rgb == "00A933":
                    bg_ok = True

        if bg_ok:
            print("✓ Background colour matches #00A933 (Green 2)")
            score += 0.5
        else:
            print("✗ Background colour is NOT #00A933")
    except Exception as exc:
        print(f"✗ Error while checking background colour: {exc}")

    # --- Requirement 2: Title text colour is pure white ---------------------
    try:
        white_ok = False

        for shape in target_slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue

            # Title placeholders have placeholder_format.type 1 (TITLE)
            # or 3 (CENTER_TITLE) in python-pptx
            if getattr(shape, "is_placeholder", False):
                ph_type = shape.placeholder_format.type
                if ph_type in (1, 3):  # Title placeholders
                    tf = shape.text_frame
                    for para in tf.paragraphs:
                        for run in para.runs:
                            font_colour = run.font.color
                            if getattr(font_colour, "type", None) == 1:  # RGB
                                rgb = str(font_colour.rgb).upper()
                                if rgb == "FFFFFF":
                                    white_ok = True
                                else:
                                    print(f"Found non-white title run colour: {rgb}")
                            else:
                                print("Title run colour not explicitly RGB or unset")

        if white_ok:
            print("✓ Title text colour is pure white (#FFFFFF)")
            score += 0.5
        else:
            print("✗ Title text colour is NOT pure white (#FFFFFF)")
    except Exception as exc:
        print(f"✗ Error while checking title text colour: {exc}")

    # --- Final scoring ------------------------------------------------------
    final_score = min(score, max_score)
    print(f"Reward score: {final_score}")
    return final_score


if __name__ == "__main__":
    FILE_PATH = (
        "/home/user/slide_295_is_the_only_one_that_hasnt_been_updated_yet_in_libreoffice_"
        "impress_how_do_i_switch_its_bac_golden.pptx"
    )

    reward_value = verify_slide_295_background_and_title(FILE_PATH)
    print(f"REWARD: {reward_value}")
