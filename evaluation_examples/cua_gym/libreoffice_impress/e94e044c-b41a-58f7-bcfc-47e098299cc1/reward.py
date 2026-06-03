"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 22 is a mess—my three shapes are crammed together on the top border. In LibreOffice Impress, how do I space those three shapes so they’re perfectly and evenly distributed across the entire top edge?
Generated: 2025-09-10 22:53:48
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation

def verify_even_distribution_top_edge(file_path: str) -> float:
    """Verify that exactly three shapes on slide 22 are evenly distributed
    across (almost) the entire top edge of the slide in a PPTX file.

    Progressive scoring (adds up to 1.0):
    1. Slide 22 exists ............................................. 0.10
    2. >=3 shapes reside very near the top edge .................... 0.10
    3. Those shapes collectively span ~100 % of the slide width .... 0.30
    4. Horizontal gaps between the three shapes are (almost) equal . 0.30
    5. Exactly three shapes were found/evaluated ................... 0.20
    """

    max_score = 1.0
    score = 0.0

    # ------------------------------------------------------------------
    # 1. Load the presentation file (no points for just loading!)
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File not found :", file_path)
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print("✗ Unable to open PPTX :", e)
        print("REWARD: 0.0")
        return 0.0

    # ------------------------------------------------------------------
    # 2. Verify slide 22 exists (index 21) -----------------------------
    # ------------------------------------------------------------------
    if len(prs.slides) >= 22:
        score += 0.10
        print("✓ Slide 22 exists (0.10)")
        slide = prs.slides[21]
    else:
        print("✗ Presentation has fewer than 22 slides")
        print(f"REWARD: {score}")
        return score

    # ------------------------------------------------------------------
    # 3. Collect shapes that sit very close to the top edge ------------
    # ------------------------------------------------------------------
    TOP_THRESHOLD = 300_000          # 0.33 inch in EMUs (36 000 EMU = 1 pt)
    top_shapes = [sh for sh in slide.shapes if hasattr(sh, "top") and sh.top <= TOP_THRESHOLD]
    print(f"Found {len(top_shapes)} shape(s) near the top edge")

    if len(top_shapes) >= 3:
        score += 0.10
        print("✓ At least three shapes located on the top edge (0.10)")
    else:
        print("✗ Fewer than three shapes on the top edge – cannot verify spacing")
        print(f"REWARD: {score}")
        return score

    # Keep only the three left-most shapes for spacing evaluation
    top_shapes = sorted(top_shapes, key=lambda sh: sh.left)[:3]

    # ------------------------------------------------------------------
    # 4. Check that the shapes span ~full slide width ------------------
    # ------------------------------------------------------------------
    slide_width = prs.slide_width
    left_most  = min(sh.left for sh in top_shapes)
    right_most = max(sh.left + sh.width for sh in top_shapes)

    EDGE_TOL = 200_000  # allow ~0.22 inch margin on both sides
    if left_most <= EDGE_TOL and right_most >= slide_width - EDGE_TOL:
        score += 0.30
        print("✓ Shapes span (almost) the full width of the slide (0.30)")
    else:
        print("✗ Shapes do NOT span the full width of the slide")

    # ------------------------------------------------------------------
    # 5. Verify equal horizontal gaps between consecutive shapes -------
    # ------------------------------------------------------------------
    shapes_sorted = sorted(top_shapes, key=lambda sh: sh.left)
    gap1 = shapes_sorted[1].left - (shapes_sorted[0].left + shapes_sorted[0].width)
    gap2 = shapes_sorted[2].left - (shapes_sorted[1].left + shapes_sorted[1].width)

    GAP_TOL = 50_000  # allow ~0.05 inch difference between gaps
    print(f"Gap 1: {gap1} EMU, Gap 2: {gap2} EMU (tol ±{GAP_TOL})")

    if gap1 > 0 and gap2 > 0 and abs(gap1 - gap2) <= GAP_TOL:
        score += 0.30
        print("✓ Gaps between shapes are (almost) equal (0.30)")
    else:
        print("✗ Gaps are uneven – shapes are not evenly distributed")

    # ------------------------------------------------------------------
    # 6. Bonus: exactly three shapes evaluated -------------------------
    # ------------------------------------------------------------------
    if len(top_shapes) == 3:
        score += 0.20
        print("✓ Exactly three shapes considered (0.20)")
    else:
        print("✗ More than three shapes at top edge – extra objects present")

    # ------------------------------------------------------------------
    # 7. Clamp & report final score ------------------------------------
    # ------------------------------------------------------------------
    final_score = min(score, max_score)
    print(f"Total score : {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score

# ----------------------------------------------------------------------
# Execute verification when run as a script -----------------------------
# ----------------------------------------------------------------------
if __name__ == "__main__":
    TEST_FILE = "/home/user/slide_22_is_a_messmy_three_shapes_are_crammed_together_on_the_top_border_in_libreoffice_impress_how__golden.pptx"
    verify_even_distribution_top_edge(TEST_FILE)

