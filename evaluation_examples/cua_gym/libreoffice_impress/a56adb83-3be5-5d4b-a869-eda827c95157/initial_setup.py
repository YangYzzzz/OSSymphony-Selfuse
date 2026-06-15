"""
Initial Setup: Export slides as PNG images
Task ID: impress_fix_069
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_069'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q3 2025 Marketing Strategy"
    slide1.placeholders[1].text = "Prepared by the Digital Marketing Team\nConfidential"
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)
    for para in slide1.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.size = Pt(40)
    for para in slide1.placeholders[1].text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Market Analysis & Competitive Landscape"
    items = [
        "Customer Segmentation Results",
        "Campaign Performance Review",
        "Social Media Growth Metrics",
        "Budget Allocation for Q4",
        "Key Milestones & Timeline",
        "Team Assignments & Next Steps",
    ]
    for item in items:
        p = body2.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 3: Market Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    tb3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Market Overview"
    p3.runs[0].font.size = Pt(32)
    p3.runs[0].font.bold = True
    p3.runs[0].font.color.rgb = RGBColor(0x1B, 0x4F, 0x72)

    # Add a table with market data
    tbl_shape = slide3.shapes.add_table(5, 4, Inches(0.8), Inches(1.8), Inches(10), Inches(3))
    tbl = tbl_shape.table
    headers = ["Region", "Market Size ($M)", "Growth Rate", "Our Share"]
    data_rows = [
        ["North America", "$245.8", "12.3%", "18.5%"],
        ["Europe", "$189.4", "8.7%", "14.2%"],
        ["Asia Pacific", "$312.1", "22.1%", "9.8%"],
        ["Latin America", "$78.6", "15.4%", "6.3%"],
    ]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row_data in enumerate(data_rows, 1):
        for c, val in enumerate(row_data):
            tbl.cell(r, c).text = val

    # --- Slide 4: Customer Segments ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    tb4_title = slide4.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1))
    tf4t = tb4_title.text_frame
    tf4t.paragraphs[0].text = "Customer Segmentation"
    tf4t.paragraphs[0].runs[0].font.size = Pt(32)
    tf4t.paragraphs[0].runs[0].font.bold = True
    tf4t.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1B, 0x4F, 0x72)

    segments = [
        ("Enterprise (42%)", "Annual contract value $50K+\nDecision cycle: 3-6 months\nKey verticals: Finance, Healthcare, Tech", RGBColor(0x2E, 0x86, 0xAB)),
        ("Mid-Market (31%)", "Annual contract value $10K-$50K\nDecision cycle: 1-3 months\nGrowing at 18% YoY", RGBColor(0xA2, 0x3B, 0x72)),
        ("SMB (27%)", "Annual contract value <$10K\nSelf-serve onboarding\nHighest volume, lowest touch", RGBColor(0xF1, 0x8F, 0x01)),
    ]
    for i, (title, desc, color) in enumerate(segments):
        left = Inches(0.8 + i * 4.0)
        box = slide4.shapes.add_textbox(left, Inches(2.0), Inches(3.5), Inches(3.5))
        tf = box.text_frame
        tf.word_wrap = True
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.runs[0].font.bold = True
        p_title.runs[0].font.size = Pt(18)
        p_title.runs[0].font.color.rgb = color
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.space_before = Pt(8)

    # --- Slide 5: Campaign Performance ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    tb5 = slide5.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1))
    tf5 = tb5.text_frame
    tf5.paragraphs[0].text = "Campaign Performance - H1 2025"
    tf5.paragraphs[0].runs[0].font.size = Pt(28)
    tf5.paragraphs[0].runs[0].font.bold = True

    tbl5_shape = slide5.shapes.add_table(6, 5, Inches(0.5), Inches(1.8), Inches(11.5), Inches(4))
    tbl5 = tbl5_shape.table
    h5 = ["Campaign", "Impressions", "Click Rate", "Conversions", "ROI"]
    data5 = [
        ["Spring Product Launch", "2.4M", "3.8%", "12,450", "340%"],
        ["Webinar Series", "890K", "5.2%", "4,230", "520%"],
        ["LinkedIn Thought Leadership", "1.1M", "2.1%", "3,890", "280%"],
        ["Email Nurture Sequence", "450K", "8.7%", "6,780", "610%"],
        ["PPC Brand Campaign", "3.2M", "1.9%", "8,920", "190%"],
    ]
    for c, h in enumerate(h5):
        tbl5.cell(0, c).text = h
        for run in tbl5.cell(0, c).text_frame.paragraphs[0].runs:
            run.font.bold = True
    for r, rd in enumerate(data5, 1):
        for c, v in enumerate(rd):
            tbl5.cell(r, c).text = v

    # --- Slide 6: Social Media Metrics ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    tb6 = slide6.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1))
    tf6 = tb6.text_frame
    tf6.paragraphs[0].text = "Social Media Growth"
    tf6.paragraphs[0].runs[0].font.size = Pt(32)
    tf6.paragraphs[0].runs[0].font.bold = True
    tf6.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1B, 0x4F, 0x72)

    metrics = [
        ("LinkedIn", "+28% followers", "142K total", Inches(1)),
        ("Twitter/X", "+15% engagement", "89K total", Inches(4.5)),
        ("Instagram", "+45% reach", "67K total", Inches(8)),
    ]
    for name, growth, total, left in metrics:
        box = slide6.shapes.add_textbox(left, Inches(2.5), Inches(3), Inches(3))
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = name
        p1.alignment = PP_ALIGN.CENTER
        p1.runs[0].font.size = Pt(24)
        p1.runs[0].font.bold = True
        p2 = tf.add_paragraph()
        p2.text = growth
        p2.alignment = PP_ALIGN.CENTER
        p2.runs[0].font.size = Pt(18)
        p2.runs[0].font.color.rgb = RGBColor(0x27, 0xAE, 0x60)
        p3 = tf.add_paragraph()
        p3.text = total
        p3.alignment = PP_ALIGN.CENTER

    # --- Slide 7: Budget Allocation ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    tb7 = slide7.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1))
    tf7 = tb7.text_frame
    tf7.paragraphs[0].text = "Q4 Budget Allocation"
    tf7.paragraphs[0].runs[0].font.size = Pt(32)
    tf7.paragraphs[0].runs[0].font.bold = True

    budget_items = [
        "Digital Advertising: $180,000 (35%)",
        "Content Production: $95,000 (18%)",
        "Events & Conferences: $72,000 (14%)",
        "Marketing Technology: $65,000 (13%)",
        "PR & Communications: $52,000 (10%)",
        "Research & Analytics: $36,000 (7%)",
        "Contingency: $15,000 (3%)",
    ]
    tb7b = slide7.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(8), Inches(4.5))
    tf7b = tb7b.text_frame
    tf7b.word_wrap = True
    tf7b.paragraphs[0].text = budget_items[0]
    tf7b.paragraphs[0].runs[0].font.size = Pt(18)
    for item in budget_items[1:]:
        p = tf7b.add_paragraph()
        p.text = item
        p.runs[0].font.size = Pt(18)
        p.space_before = Pt(8)

    # --- Slide 8: Next Steps & Timeline ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    tb8 = slide8.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1))
    tf8 = tb8.text_frame
    tf8.paragraphs[0].text = "Next Steps & Timeline"
    tf8.paragraphs[0].runs[0].font.size = Pt(32)
    tf8.paragraphs[0].runs[0].font.bold = True
    tf8.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1B, 0x4F, 0x72)

    steps = [
        ("Week 1-2", "Finalize campaign creative assets and landing pages"),
        ("Week 3-4", "Launch LinkedIn and Google Ads campaigns"),
        ("Week 5-6", "Execute webinar series with industry partners"),
        ("Week 7-8", "Mid-quarter performance review and optimization"),
        ("Week 9-10", "Scale high-performing channels, pause underperformers"),
        ("Week 11-12", "Final reporting, Q1 2026 planning kickoff"),
    ]
    for i, (week, desc) in enumerate(steps):
        y = Inches(1.8 + i * 0.85)
        tb_w = slide8.shapes.add_textbox(Inches(1.0), y, Inches(2.0), Inches(0.7))
        tf_w = tb_w.text_frame
        tf_w.paragraphs[0].text = week
        tf_w.paragraphs[0].runs[0].font.bold = True
        tf_w.paragraphs[0].runs[0].font.size = Pt(16)
        tf_w.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x2E, 0x86, 0xAB)
        tb_d = slide8.shapes.add_textbox(Inches(3.5), y, Inches(8), Inches(0.7))
        tf_d = tb_d.text_frame
        tf_d.paragraphs[0].text = desc
        tf_d.paragraphs[0].runs[0].font.size = Pt(16)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Ensure ~/Desktop/slide_images/ does NOT exist
    subprocess.run(['rm', '-rf', f'{WORKDIR}/Desktop/slide_images'], check=False)

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
