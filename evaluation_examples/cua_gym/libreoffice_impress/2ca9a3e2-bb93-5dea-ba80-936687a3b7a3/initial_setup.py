"""
Initial Setup: 5-slide product comparison presentation
Task ID: osworld_impress_underline_darkred_table_006
Domain: libreoffice_impress

Slide 2 has two body textboxes and a 5x3 comparison table, all in black text.
The task requires applying underline + #DC143C to all slide 2 content.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_underline_darkred_table_006'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_text_black(tf, text_lines):
    """Set text frame content with black, non-underlined text."""
    # Clear existing paragraphs beyond first
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)

    for idx, line in enumerate(text_lines):
        if idx == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.text = line
        for run in para.runs:
            run.font.size = Pt(14)
            run.font.bold = False
            run.font.italic = False
            run.font.underline = False
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    BLACK = RGBColor(0x00, 0x00, 0x00)

    # ----------------------------------------------------------------
    # Slide 1: Title slide
    # ----------------------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "TechVision Pro: Product Comparison 2025"
    slide1.placeholders[1].text = "Helping you choose the right solution"
    for para in slide1.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = BLACK

    # ----------------------------------------------------------------
    # Slide 2: Comparison overview (two textboxes + 5x3 table)
    # ----------------------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    # Slide 2 title
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.paragraphs[0].text = "Product Overview"
    for run in title_tf.paragraphs[0].runs:
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.underline = False
        run.font.color.rgb = BLACK

    # Textbox 1: Product highlights
    tb1 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.2), Inches(1.6))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    lines1 = [
        "Key Highlights",
        "- Industry-leading performance metrics",
        "- Certified for enterprise deployments",
        "- 24/7 dedicated support included",
    ]
    for idx, line in enumerate(lines1):
        if idx == 0:
            para = tf1.paragraphs[0]
        else:
            para = tf1.add_paragraph()
        para.text = line
        for run in para.runs:
            run.font.size = Pt(13)
            run.font.bold = False
            run.font.underline = False
            run.font.color.rgb = BLACK

    # Textbox 2: Availability note
    tb2 = slide2.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.2), Inches(1.6))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    lines2 = [
        "Availability",
        "- Global distribution in 42 countries",
        "- Immediate licensing available",
        "- Free 30-day evaluation period",
    ]
    for idx, line in enumerate(lines2):
        if idx == 0:
            para = tf2.paragraphs[0]
        else:
            para = tf2.add_paragraph()
        para.text = line
        for run in para.runs:
            run.font.size = Pt(13)
            run.font.bold = False
            run.font.underline = False
            run.font.color.rgb = BLACK

    # 5x3 comparison table on slide 2
    tbl_shape = slide2.shapes.add_table(5, 3, Inches(0.5), Inches(3.0), Inches(9), Inches(3.8))
    tbl = tbl_shape.table

    # Column widths
    tbl.columns[0].width = Inches(3.5)
    tbl.columns[1].width = Inches(2.75)
    tbl.columns[2].width = Inches(2.75)

    table_data = [
        ["Feature",          "Standard Edition",    "Professional Edition"],
        ["Storage Capacity", "500 GB SSD",           "2 TB NVMe SSD"],
        ["CPU Cores",        "4 cores / 8 threads",  "16 cores / 32 threads"],
        ["Annual License",   "$1,200 / year",        "$4,800 / year"],
        ["Support Level",    "Business hours only",  "24/7 Premium support"],
    ]

    for r_idx, row_data in enumerate(table_data):
        for c_idx, cell_text in enumerate(row_data):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = cell_text
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(13)
                    run.font.bold = (r_idx == 0)
                    run.font.underline = False
                    run.font.color.rgb = BLACK

    # ----------------------------------------------------------------
    # Slide 3: Performance benchmarks
    # ----------------------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Performance Benchmarks"
    slide3.placeholders[1].text = (
        "Our products consistently outperform industry averages.\n"
        "Benchmark tests conducted Q3 2024 by independent lab.\n"
        "Standard Edition: 2,400 IOPS average throughput.\n"
        "Professional Edition: 18,500 IOPS average throughput.\n"
        "Both editions certified under ISO/IEC 27001 standards."
    )
    for para in slide3.placeholders[1].text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(14)
            run.font.color.rgb = BLACK

    # ----------------------------------------------------------------
    # Slide 4: Customer testimonials
    # ----------------------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Customer Testimonials"
    slide4.placeholders[1].text = (
        "\"The Professional Edition transformed our workflow.\" — Aria Nguyen, CTO, NexaTech\n"
        "\"Setup was seamless and support is world-class.\" — Marcus Ellison, VP Engineering\n"
        "\"ROI was visible within the first quarter.\" — Sofia Delgado, Operations Director\n"
        "\"Best enterprise purchase we made this year.\" — James Harrington, IT Manager"
    )
    for para in slide4.placeholders[1].text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(14)
            run.font.color.rgb = BLACK

    # ----------------------------------------------------------------
    # Slide 5: Next steps
    # ----------------------------------------------------------------
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Next Steps"
    slide5.placeholders[1].text = (
        "1. Request a personalized demo from your account manager.\n"
        "2. Review the technical specifications document.\n"
        "3. Begin a 30-day free evaluation trial.\n"
        "4. Contact procurement for volume licensing options.\n"
        "5. Join our upcoming webinar: Q1 Product Roadmap Preview."
    )
    for para in slide5.placeholders[1].text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(14)
            run.font.color.rgb = BLACK

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
