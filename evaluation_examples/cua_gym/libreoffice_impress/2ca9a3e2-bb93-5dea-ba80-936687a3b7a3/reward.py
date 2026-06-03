"""
Reward Script: Apply underline + dark crimson color (#DC143C) to ALL content on slide 2
Task ID: osworld_impress_underline_darkred_table_006
Domain: libreoffice_impress
Scoring:
  Component 1: All textbox runs on slide 2 have underline=True          (0.25 pts)
  Component 2: All textbox runs on slide 2 have color=#DC143C           (0.25 pts)
  Component 3: All table cell runs on slide 2 have underline=True       (0.25 pts)
  Component 4: All table cell runs on slide 2 have color=#DC143C        (0.25 pts)
  Total: 1.0
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_underline_darkred_table_006'
TARGET_COLOR = 'DC143C'


def verify_task(file_path):
    """
    Verify that all content on slide 2 (index 1) has been given underline formatting
    and dark crimson color (#DC143C). Checks textboxes and table cells independently.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"CRITICAL: Presentation has only {len(prs.slides)} slides, expected >= 2")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[1]  # Slide 2 (0-indexed)

    # ----------------------------------------------------------------
    # Component 1: All TEXTBOX runs on slide 2 have underline=True (0.25 pts)
    # ----------------------------------------------------------------
    try:
        textbox_runs_found = 0
        textbox_underline_pass = 0
        textbox_underline_fails = []

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                continue  # tables handled separately
            if not shape.has_text_frame:
                continue
            # Only check TEXT_BOX shapes (type 17), skip title placeholder
            if shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
                continue

            for para in shape.text_frame.paragraphs:
                non_empty_runs = [r for r in para.runs if (r.text or '').strip()]
                for run in non_empty_runs:
                    textbox_runs_found += 1
                    if run.font.underline is True:
                        textbox_underline_pass += 1
                    else:
                        textbox_underline_fails.append(
                            f"shape={shape.name!r}, text={run.text!r}, underline={run.font.underline}"
                        )

        if textbox_runs_found == 0:
            print("FAIL: Component 1 — No textbox runs found on slide 2")
        elif textbox_underline_pass == textbox_runs_found:
            print(f"PASS: Component 1 — All {textbox_runs_found} textbox runs have underline=True (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — {textbox_underline_pass}/{textbox_runs_found} textbox runs have underline=True")
            for detail in textbox_underline_fails[:5]:
                print(f"  Missing underline: {detail}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # ----------------------------------------------------------------
    # Component 2: All TEXTBOX runs on slide 2 have color=#DC143C (0.25 pts)
    # ----------------------------------------------------------------
    try:
        textbox_color_pass = 0
        textbox_color_fails = []
        textbox_color_runs_found = 0

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                continue
            if not shape.has_text_frame:
                continue
            if shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
                continue

            for para in shape.text_frame.paragraphs:
                non_empty_runs = [r for r in para.runs if (r.text or '').strip()]
                for run in non_empty_runs:
                    textbox_color_runs_found += 1
                    try:
                        if run.font.color.type is not None:
                            actual_color = str(run.font.color.rgb).upper()
                            if actual_color == TARGET_COLOR:
                                textbox_color_pass += 1
                            else:
                                textbox_color_fails.append(
                                    f"shape={shape.name!r}, text={run.text!r}, color={actual_color}"
                                )
                        else:
                            textbox_color_fails.append(
                                f"shape={shape.name!r}, text={run.text!r}, color=None (inherited)"
                            )
                    except Exception as ce:
                        textbox_color_fails.append(
                            f"shape={shape.name!r}, text={run.text!r}, color_error={ce}"
                        )

        if textbox_color_runs_found == 0:
            print("FAIL: Component 2 — No textbox runs found on slide 2")
        elif textbox_color_pass == textbox_color_runs_found:
            print(f"PASS: Component 2 — All {textbox_color_runs_found} textbox runs have color=#{TARGET_COLOR} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 — {textbox_color_pass}/{textbox_color_runs_found} textbox runs have color=#{TARGET_COLOR}")
            for detail in textbox_color_fails[:5]:
                print(f"  Wrong color: {detail}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # ----------------------------------------------------------------
    # Component 3: All TABLE cell runs on slide 2 have underline=True (0.25 pts)
    # ----------------------------------------------------------------
    try:
        table_shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE]
        table_runs_found = 0
        table_underline_pass = 0
        table_underline_fails = []

        for shape in table_shapes:
            table = shape.table
            num_rows = len(table.rows)
            num_cols = len(table.columns)

            for r in range(num_rows):
                for c in range(num_cols):
                    cell = table.cell(r, c)
                    for para in cell.text_frame.paragraphs:
                        non_empty_runs = [run for run in para.runs if (run.text or '').strip()]
                        for run in non_empty_runs:
                            table_runs_found += 1
                            if run.font.underline is True:
                                table_underline_pass += 1
                            else:
                                table_underline_fails.append(
                                    f"cell({r},{c}), text={run.text!r}, underline={run.font.underline}"
                                )

        if len(table_shapes) == 0:
            print("FAIL: Component 3 — No table found on slide 2")
        elif table_runs_found == 0:
            print("FAIL: Component 3 — Table found but no runs detected")
        elif table_underline_pass == table_runs_found:
            print(f"PASS: Component 3 — All {table_runs_found} table cell runs have underline=True (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 — {table_underline_pass}/{table_runs_found} table cell runs have underline=True")
            for detail in table_underline_fails[:5]:
                print(f"  Missing underline: {detail}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # ----------------------------------------------------------------
    # Component 4: All TABLE cell runs on slide 2 have color=#DC143C (0.25 pts)
    # ----------------------------------------------------------------
    try:
        table_color_pass = 0
        table_color_fails = []
        table_color_runs_found = 0

        for shape in [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE]:
            table = shape.table
            num_rows = len(table.rows)
            num_cols = len(table.columns)

            for r in range(num_rows):
                for c in range(num_cols):
                    cell = table.cell(r, c)
                    for para in cell.text_frame.paragraphs:
                        non_empty_runs = [run for run in para.runs if (run.text or '').strip()]
                        for run in non_empty_runs:
                            table_color_runs_found += 1
                            try:
                                if run.font.color.type is not None:
                                    actual_color = str(run.font.color.rgb).upper()
                                    if actual_color == TARGET_COLOR:
                                        table_color_pass += 1
                                    else:
                                        table_color_fails.append(
                                            f"cell({r},{c}), text={run.text!r}, color={actual_color}"
                                        )
                                else:
                                    table_color_fails.append(
                                        f"cell({r},{c}), text={run.text!r}, color=None (inherited)"
                                    )
                            except Exception as ce:
                                table_color_fails.append(
                                    f"cell({r},{c}), text={run.text!r}, color_error={ce}"
                                )

        if table_color_runs_found == 0:
            print("FAIL: Component 4 — Table found but no runs detected for color check")
        elif table_color_pass == table_color_runs_found:
            print(f"PASS: Component 4 — All {table_color_runs_found} table cell runs have color=#{TARGET_COLOR} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 4 — {table_color_pass}/{table_color_runs_found} table cell runs have color=#{TARGET_COLOR}")
            for detail in table_color_fails[:5]:
                print(f"  Wrong color: {detail}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
