"""
Initial Setup: Art History lecture presentation with 8 slides
Task ID: impress_teach_014
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_014'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_content_slide(prs, layout_idx, title_text, body_lines):
    """Add a slide with a title and bullet-point body content."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title_text
    # Use placeholder 1 for body content
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Introduction to Art History"
    slide1.placeholders[1].text = "A Comprehensive Survey of Western Art\nProfessor Elena Vasquez\nFall 2025"

    # --- Slide 2: Renaissance Masters ---
    add_content_slide(prs, 1, "Renaissance Masters", [
        "Leonardo da Vinci (1452-1519) - Mona Lisa, The Last Supper",
        "Michelangelo Buonarroti (1475-1564) - Sistine Chapel ceiling",
        "Raphael Sanzio (1483-1520) - The School of Athens",
        "Sandro Botticelli (1445-1510) - The Birth of Venus",
        "Key characteristics: humanism, perspective, naturalism",
    ])

    # --- Slide 3: Baroque Period ---
    add_content_slide(prs, 1, "Baroque Period", [
        "Caravaggio (1571-1610) - dramatic use of chiaroscuro",
        "Peter Paul Rubens (1577-1640) - dynamic compositions",
        "Rembrandt van Rijn (1606-1669) - masterful portraiture",
        "Gian Lorenzo Bernini (1598-1680) - sculptural emotion",
        "Key characteristics: grandeur, movement, contrast of light and shadow",
    ])

    # --- Slide 4: Impressionism Movement ---
    add_content_slide(prs, 1, "Impressionism Movement", [
        "Claude Monet (1840-1926) - Water Lilies, Impression Sunrise",
        "Pierre-Auguste Renoir (1841-1919) - luminous figure paintings",
        "Edgar Degas (1834-1917) - ballet and movement studies",
        "Berthe Morisot (1841-1895) - intimate domestic scenes",
        "Key characteristics: light, color, visible brushstrokes, everyday subjects",
    ])

    # --- Slide 5: Lecture Overview (THIS IS THE SLIDE TO BE MOVED) ---
    add_content_slide(prs, 1, "Lecture Overview", [
        "1. Renaissance Masters - The rebirth of classical ideals",
        "2. Baroque Period - Drama and emotional intensity",
        "3. Impressionism - Capturing light and modern life",
        "4. Modern Art - Breaking with tradition",
        "5. Contemporary Trends - Art in the 21st century",
        "6. Discussion and critical analysis",
    ])

    # --- Slide 6: Modern Art Developments ---
    add_content_slide(prs, 1, "Modern Art Developments", [
        "Pablo Picasso (1881-1973) - Cubism and beyond",
        "Wassily Kandinsky (1866-1944) - Abstract Expressionism pioneer",
        "Frida Kahlo (1907-1954) - Surrealism and self-portraiture",
        "Jackson Pollock (1912-1956) - Action painting",
        "Key characteristics: abstraction, experimentation, challenging norms",
    ])

    # --- Slide 7: Contemporary Trends ---
    add_content_slide(prs, 1, "Contemporary Trends", [
        "Banksy - Street art as political commentary",
        "Yayoi Kusama - Immersive installations and infinity rooms",
        "Ai Weiwei - Art as activism and social critique",
        "Kara Walker - Silhouettes exploring race and identity",
        "Key characteristics: digital media, installation, social engagement",
    ])

    # --- Slide 8: Discussion & Questions ---
    add_content_slide(prs, 1, "Discussion & Questions", [
        "How did the Renaissance set the stage for modern art?",
        "What role does cultural context play in artistic movements?",
        "Compare the emotional impact of Baroque vs. Impressionism",
        "How does contemporary art challenge our definitions of art?",
        "Next lecture: Post-Impressionism and the road to abstraction",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
