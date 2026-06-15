"""
Reward Script: Reorder slides — move slide 5 ("Lecture Overview") to become slide 2.
Task ID: impress_teach_014
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Slide 2 title is "Lecture Overview"
  Component 2 (0.3): Total slide count remains 8
  Component 3 (0.3): Full slide order matches expected sequence after the move
"""

import os

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_014'

# Expected slide order after moving original slide 5 to position 2
EXPECTED_TITLES = [
    "Introduction to Art History",
    "Lecture Overview",
    "Renaissance Masters",
    "Baroque Period",
    "Impressionism Movement",
    "Modern Art Developments",
    "Contemporary Trends",
    "Discussion & Questions",
]


def get_slide_first_text(slide):
    """Extract the first non-empty text from a slide's shapes."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    return t
    return ""


def persist_app_state():
    """Best-effort save via Ctrl+S in case file is open in LibreOffice."""
    try:
        os.environ["DISPLAY"] = ":0"
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        import time
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)
    actual_titles = [get_slide_first_text(s) for s in slides]
    print(f"INFO: Found {len(slides)} slides with titles: {actual_titles}")

    # Component 1: Slide 2 title is "Lecture Overview" (0.4 points)
    # This is the core task requirement — the moved slide must be in position 2.
    # INITIAL: Slide 2 = "Renaissance Masters" -> FAIL
    # GOLDEN:  Slide 2 = "Lecture Overview" -> PASS
    try:
        if len(slides) >= 2:
            slide2_title = actual_titles[1]  # 0-indexed
            if slide2_title == "Lecture Overview":
                print(f"PASS: Component 1 — Slide 2 title is 'Lecture Overview' (0.4 pts)")
                total_score += 0.4
            else:
                print(f"FAIL: Component 1 — Slide 2 title is '{slide2_title}', expected 'Lecture Overview'")
        else:
            print(f"FAIL: Component 1 — Not enough slides (found {len(slides)})")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Total slide count remains 8 (0.3 points)
    # The move should not add or remove any slides.
    # IMPORTANT: This is scored as a compound check — count is 8 AND slide 2 is correct.
    # In initial_env, slide 2 is NOT "Lecture Overview", so the compound check fails.
    try:
        if len(slides) == 8 and (len(slides) >= 2 and actual_titles[1] == "Lecture Overview"):
            print(f"PASS: Component 2 — Slide count is 8 and slide 2 is correct (0.3 pts)")
            total_score += 0.3
        else:
            if len(slides) != 8:
                print(f"FAIL: Component 2 — Slide count is {len(slides)}, expected 8")
            else:
                print(f"FAIL: Component 2 — Slide count OK but slide 2 not yet moved")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Full slide order matches expected sequence (0.3 points)
    # All 8 slides should be in the correct order after the move.
    # INITIAL: Order has "Lecture Overview" at position 5 -> FAIL
    # GOLDEN:  Order has "Lecture Overview" at position 2 -> PASS
    try:
        if actual_titles == EXPECTED_TITLES:
            print(f"PASS: Component 3 — Full slide order matches expected sequence (0.3 pts)")
            total_score += 0.3
        else:
            mismatches = []
            for i, (exp, act) in enumerate(zip(EXPECTED_TITLES, actual_titles)):
                if exp != act:
                    mismatches.append(f"Slide {i+1}: expected '{exp}', found '{act}'")
            if len(actual_titles) != len(EXPECTED_TITLES):
                mismatches.append(f"Count mismatch: expected {len(EXPECTED_TITLES)}, found {len(actual_titles)}")
            print(f"FAIL: Component 3 — Slide order mismatches: {'; '.join(mismatches)}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
