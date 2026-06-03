"""
Reward Script: Reorder slides — move Conclusion from slide 7 to slide 10
Task ID: impress_stu_011
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Slide 7 title contains "Appendix A" (was "Conclusion" initially)
  Component 2 (0.20): Slide 8 title contains "Appendix B"
  Component 3 (0.20): Slide 9 title contains "Appendix C" / "Future Work"
  Component 4 (0.30): Slide 10 title contains "Conclusion"
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_011'


def get_slide_title(slide):
    """Extract first non-empty text from a slide as its title."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    return t
    return ""


def persist_app_state():
    """Send Ctrl+S to save any unsaved LibreOffice changes."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0

    The task reorders slides so that:
      - Original slide 7 (Conclusion) moves to slide 10
      - Original slides 8, 9, 10 shift up to positions 7, 8, 9
    We verify by checking the title/content of slides 7-10.
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must still have 10 slides
    num_slides = len(prs.slides)
    if num_slides != 10:
        print(f"FAIL: Precondition — expected 10 slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: slides 1-6 should be unchanged (gate, not scored)
    expected_first_six = [
        "Capstone Project",
        "Project Overview",
        "Methodology",
        "Key Findings",
        "Dashboard Demo",
        "Policy Recommendations",
    ]
    for i in range(6):
        title = get_slide_title(prs.slides[i])
        if expected_first_six[i].lower() not in title.lower():
            print(f"FAIL: Precondition — slide {i+1} title '{title}' does not contain '{expected_first_six[i]}'")
            print("REWARD: 0.0")
            return 0.0
    print("PASS: Precondition — 10 slides, first 6 unchanged")

    # Component 1: Slide 7 should be "Appendix A: Rideshare Analysis" (0.30 points)
    # In initial state, slide 7 is "Conclusion" — this check FAILS on initial
    try:
        title_7 = get_slide_title(prs.slides[6])
        if "appendix a" in title_7.lower():
            print(f"PASS: Component 1 — Slide 7 is '{title_7}' (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 1 — Slide 7 expected 'Appendix A...', found '{title_7}'")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 8 should be "Appendix B: Transit Performance Metrics" (0.20 points)
    # In initial state, slide 8 is "Appendix A" — this check FAILS on initial
    try:
        title_8 = get_slide_title(prs.slides[7])
        if "appendix b" in title_8.lower():
            print(f"PASS: Component 2 — Slide 8 is '{title_8}' (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 2 — Slide 8 expected 'Appendix B...', found '{title_8}'")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 9 should be "Appendix C: Future Work" (0.20 points)
    # In initial state, slide 9 is "Appendix B" — this check FAILS on initial
    try:
        title_9 = get_slide_title(prs.slides[8])
        if "appendix c" in title_9.lower() or "future work" in title_9.lower():
            print(f"PASS: Component 3 — Slide 9 is '{title_9}' (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 3 — Slide 9 expected 'Appendix C...' or 'Future Work', found '{title_9}'")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 10 should be "Conclusion" (0.30 points)
    # In initial state, slide 10 is "Appendix C" — this check FAILS on initial
    try:
        title_10 = get_slide_title(prs.slides[9])
        if "conclusion" in title_10.lower():
            print(f"PASS: Component 4 — Slide 10 is '{title_10}' (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 4 — Slide 10 expected 'Conclusion', found '{title_10}'")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
