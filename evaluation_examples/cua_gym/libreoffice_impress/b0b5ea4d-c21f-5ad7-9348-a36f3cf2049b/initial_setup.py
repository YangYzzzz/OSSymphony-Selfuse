"""
Initial Setup: Create a 10-slide Capstone Presentation with slide 7 as 'Conclusion'
Task ID: impress_stu_011
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_011'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.paragraphs[0].text = point
        else:
            p = tf.add_paragraph()
            p.text = point
    return slide


def create_initial():
    prs = Presentation()

    # Slide 1: Title Slide
    add_title_slide(
        prs,
        "Capstone Project: Urban Mobility Analytics",
        "Department of Data Science\nPresented by: Anika Sharma & Leo Park\nDecember 2025"
    )

    # Slide 2: Project Overview
    add_content_slide(prs, "Project Overview", [
        "Objective: Analyze urban transportation patterns in metro areas",
        "Data sources: GPS traces, transit schedules, rideshare logs",
        "Time period: January 2024 - November 2025",
        "Partner organizations: Metro Transit Authority, CityBike",
        "Deliverables: Interactive dashboard and policy recommendations"
    ])

    # Slide 3: Methodology
    add_content_slide(prs, "Methodology", [
        "Phase 1: Data collection and cleaning (3 months)",
        "Phase 2: Exploratory analysis with spatial clustering",
        "Phase 3: Predictive modeling using gradient-boosted trees",
        "Phase 4: Dashboard development with Plotly Dash",
        "Validation: Cross-validated against 2025 Q3 actuals"
    ])

    # Slide 4: Key Findings
    add_content_slide(prs, "Key Findings", [
        "Peak commute times shifted 22 minutes later post-2024",
        "Rideshare usage up 34% in suburban corridors",
        "Bus route optimization could save $2.1M annually",
        "Bike-share stations underutilized in 7 of 15 zones",
        "Weather impact on mode choice: 18% shift in rainy conditions"
    ])

    # Slide 5: Dashboard Demo
    add_content_slide(prs, "Dashboard Demo", [
        "Real-time heatmap of transit density by hour",
        "Route efficiency scoring with color-coded overlays",
        "Demand forecasting module (7-day rolling predictions)",
        "User filter controls: mode, time range, geography",
        "Export functionality for stakeholder reports"
    ])

    # Slide 6: Recommendations
    add_content_slide(prs, "Policy Recommendations", [
        "Reallocate 3 bus routes to high-demand suburban corridors",
        "Expand bike-share stations in zones 4, 9, and 12",
        "Introduce dynamic pricing for peak-hour rideshare trips",
        "Pilot a micro-transit service in underserved neighborhoods",
        "Invest in real-time passenger information displays"
    ])

    # Slide 7: Conclusion (this is the one to be moved)
    add_content_slide(prs, "Conclusion", [
        "Urban mobility is evolving rapidly with new data sources",
        "Our model achieves 89% accuracy on route demand prediction",
        "Dashboard adopted by Metro Transit for Q1 2026 planning",
        "Collaboration between agencies is key to implementation",
        "Thank you for your attention - Questions welcome"
    ])

    # Slide 8: Additional Data - Rideshare Analysis
    add_content_slide(prs, "Appendix A: Rideshare Analysis", [
        "Total trips analyzed: 4.2 million across 18 months",
        "Average trip distance: 5.8 km (urban), 14.3 km (suburban)",
        "Peak hours: 7:30-9:00 AM and 5:00-6:30 PM",
        "Driver utilization rate: 73% during peak, 41% off-peak",
        "Revenue per trip: $12.40 average, $18.90 suburban"
    ])

    # Slide 9: Additional Data - Transit Performance
    add_content_slide(prs, "Appendix B: Transit Performance Metrics", [
        "On-time performance: 82% (target: 90%)",
        "Average passenger load factor: 67% during peak",
        "Route coverage gap: 15% of residential areas underserved",
        "Maintenance cost per vehicle-km: $1.24 (down from $1.58)",
        "Customer satisfaction survey: 3.7/5.0 overall rating"
    ])

    # Slide 10: Additional Data - Future Work
    add_content_slide(prs, "Appendix C: Future Work", [
        "Integrate electric vehicle charging station data",
        "Expand model to include freight and delivery logistics",
        "Develop mobile app for real-time commuter guidance",
        "Partner with neighboring cities for regional analysis",
        "Apply reinforcement learning for adaptive signal control"
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
