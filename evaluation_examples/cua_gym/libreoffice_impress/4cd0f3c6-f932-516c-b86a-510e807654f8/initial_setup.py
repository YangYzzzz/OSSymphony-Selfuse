"""
Initial Setup: Create a 9-slide sustainability presentation with 24pt black titles
Task ID: impress_stu_008
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_008'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide data: (title, body_text)
    slides_data = [
        (
            "Campus Sustainability Initiative 2025",
            "Presented by the Office of Environmental Stewardship\n"
            "University of Pacific Northwest\nMarch 2025"
        ),
        (
            "Our Environmental Mission",
            "We are committed to reducing our campus carbon footprint by 40% before 2030.\n"
            "This initiative aligns with the UN Sustainable Development Goals and our\n"
            "institutional responsibility to future generations."
        ),
        (
            "Current Energy Consumption",
            "Annual electricity usage: 48.7 million kWh\n"
            "Natural gas consumption: 12.3 million therms\n"
            "Campus fleet fuel: 285,000 gallons/year\n"
            "Total carbon emissions: 127,400 metric tons CO2e"
        ),
        (
            "Renewable Energy Transition",
            "Solar panel installation across 14 buildings: 3.2 MW capacity\n"
            "Wind turbine pilot program on North Campus\n"
            "Geothermal heating for the new Science Complex\n"
            "Power purchase agreement with Green Valley Wind Farm"
        ),
        (
            "Waste Reduction Progress",
            "Composting program diverted 1,200 tons from landfill in 2024\n"
            "Single-use plastic ban implemented across all dining facilities\n"
            "E-waste recycling partnership with TechCycle Solutions\n"
            "Zero-waste events policy adopted for all campus gatherings"
        ),
        (
            "Water Conservation Measures",
            "Smart irrigation systems saved 18 million gallons last year\n"
            "Low-flow fixtures installed in 92% of campus buildings\n"
            "Rainwater harvesting tanks at Recreation Center and Library\n"
            "Greywater recycling pilot in the new Residence Hall complex"
        ),
        (
            "Transportation and Mobility",
            "Electric shuttle fleet replacing diesel buses by Fall 2025\n"
            "Bike-share program expanded to 450 stations across campus\n"
            "Remote work policy reduced commuter trips by 23%\n"
            "EV charging stations: 180 installed, 120 more planned"
        ),
        (
            "Student and Faculty Engagement",
            "Green Ambassador program: 340 trained student volunteers\n"
            "Sustainability course requirement added to core curriculum\n"
            "Annual Eco-Challenge competition with $15,000 in prizes\n"
            "Faculty research grants: $2.4 million allocated for green innovation"
        ),
        (
            "Roadmap and Next Steps",
            "Phase 1 (2025): Complete solar installations and fleet electrification\n"
            "Phase 2 (2026): Achieve 50% renewable energy target\n"
            "Phase 3 (2027): Launch carbon offset marketplace\n"
            "Phase 4 (2028-2030): Reach net-zero campus operations"
        ),
    ]

    for i, (title_text, body_text) in enumerate(slides_data):
        if i == 0:
            slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content

        # Set title with 24pt, black, NOT bold
        title_shape = slide.shapes.title
        title_shape.text = ""
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title_text
        run.font.size = Pt(24)
        run.font.bold = False
        run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

        # Set body content
        if i == 0:
            # Subtitle placeholder
            subtitle = slide.placeholders[1]
            subtitle.text = body_text
        else:
            # Content placeholder
            body_shape = slide.placeholders[1]
            body_shape.text = ""
            btf = body_shape.text_frame
            lines = body_text.split('\n')
            for j, line in enumerate(lines):
                if j == 0:
                    bp = btf.paragraphs[0]
                else:
                    bp = btf.add_paragraph()
                bp.text = line.strip()

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
