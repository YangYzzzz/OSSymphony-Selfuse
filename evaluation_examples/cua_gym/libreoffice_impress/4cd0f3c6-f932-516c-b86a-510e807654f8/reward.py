"""
Reward Script: Set all slide titles to 32pt, bold, dark green (#006747)
Task ID: impress_stu_008
Domain: libreoffice_impress
Scoring:
  Component 1 (0.35): All 9 titles have font size 32pt
  Component 2 (0.30): All 9 titles are bold
  Component 3 (0.35): All 9 titles have color #006747
"""

import os

from pptx import Presentation
from pptx.util import Pt

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_008'
EXPECTED_SLIDES = 9
EXPECTED_SIZE_PT = 32
EXPECTED_BOLD = "yes"  # bold expected
EXPECTED_COLOR = '006747'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)
    num_slides = len(slides)
    if num_slides == 0:
        print("FAIL: Presentation has no slides")
        print("REWARD: 0.0")
        return 0.0

    # Collect title run properties for each slide
    title_props = []  # list of dicts: {size_pt, bold, color_hex}
    for i, slide in enumerate(slides):
        title_shape = slide.shapes.title
        if title_shape is None or not title_shape.has_text_frame:
            print(f"WARNING: Slide {i+1} has no title shape")
            title_props.append(None)
            continue

        # Gather properties from all runs in the title
        runs = []
        for para in title_shape.text_frame.paragraphs:
            for run in para.runs:
                if not (run.text or "").strip():
                    continue
                size_pt = None
                if run.font.size is not None:
                    size_pt = run.font.size / 12700  # EMU to pt

                bold = run.font.bold
                # None means inherit — treat as False for comparison
                if bold is None:
                    bold = False

                color_hex = None
                try:
                    if run.font.color.type is not None:
                        color_hex = str(run.font.color.rgb)
                except Exception:
                    pass

                runs.append({
                    'size_pt': size_pt,
                    'bold': bold,
                    'color_hex': color_hex,
                })

        if runs:
            title_props.append(runs)
        else:
            print(f"WARNING: Slide {i+1} title has no non-empty runs")
            title_props.append(None)

    # Count slides with valid title data
    valid_slides = [p for p in title_props if p is not None]
    if not valid_slides:
        print("FAIL: No slides have valid title data")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Font size = 32pt for all titles (0.35 points)
    # Award proportional credit per slide
    try:
        size_pass_count = 0
        for i, props in enumerate(title_props):
            if props is None:
                continue
            all_correct = all(
                r['size_pt'] is not None and abs(r['size_pt'] - EXPECTED_SIZE_PT) < 0.5
                for r in props
            )
            if all_correct:
                size_pass_count += 1
            else:
                actual_sizes = [r['size_pt'] for r in props]
                print(f"FAIL: Slide {i+1} title size(s): {actual_sizes}, expected {EXPECTED_SIZE_PT}pt")

        size_ratio = size_pass_count / EXPECTED_SLIDES
        size_score = 0.35 * size_ratio
        if size_pass_count == EXPECTED_SLIDES:
            print(f"PASS: Component 1 -- All {EXPECTED_SLIDES} titles are {EXPECTED_SIZE_PT}pt ({size_score:.2f} pts)")
            total_score += size_score
        elif size_pass_count > 0:
            print(f"PARTIAL: Component 1 -- {size_pass_count}/{EXPECTED_SLIDES} titles are {EXPECTED_SIZE_PT}pt ({size_score:.2f} pts)")
            total_score += size_score
        else:
            print(f"FAIL: Component 1 -- 0/{EXPECTED_SLIDES} titles are {EXPECTED_SIZE_PT}pt")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Bold = True for all titles (0.30 points)
    try:
        bold_pass_count = 0
        for i, props in enumerate(title_props):
            if props is None:
                continue
            all_bold = all(r['bold'] is True for r in props)
            if all_bold:
                bold_pass_count += 1
            else:
                actual_bolds = [r['bold'] for r in props]
                print(f"FAIL: Slide {i+1} title bold: {actual_bolds}, expected True")

        bold_ratio = bold_pass_count / EXPECTED_SLIDES
        bold_score = 0.30 * bold_ratio
        if bold_pass_count == EXPECTED_SLIDES:
            print(f"PASS: Component 2 -- All {EXPECTED_SLIDES} titles are bold ({bold_score:.2f} pts)")
            total_score += bold_score
        elif bold_pass_count > 0:
            print(f"PARTIAL: Component 2 -- {bold_pass_count}/{EXPECTED_SLIDES} titles are bold ({bold_score:.2f} pts)")
            total_score += bold_score
        else:
            print(f"FAIL: Component 2 -- 0/{EXPECTED_SLIDES} titles are bold")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Color = #006747 for all titles (0.35 points)
    try:
        color_pass_count = 0
        for i, props in enumerate(title_props):
            if props is None:
                continue
            all_color = all(
                r['color_hex'] is not None and r['color_hex'].upper() == EXPECTED_COLOR.upper()
                for r in props
            )
            if all_color:
                color_pass_count += 1
            else:
                actual_colors = [r['color_hex'] for r in props]
                print(f"FAIL: Slide {i+1} title color: {actual_colors}, expected #{EXPECTED_COLOR}")

        color_ratio = color_pass_count / EXPECTED_SLIDES
        color_score = 0.35 * color_ratio
        if color_pass_count == EXPECTED_SLIDES:
            print(f"PASS: Component 3 -- All {EXPECTED_SLIDES} titles are #{EXPECTED_COLOR} ({color_score:.2f} pts)")
            total_score += color_score
        elif color_pass_count > 0:
            print(f"PARTIAL: Component 3 -- {color_pass_count}/{EXPECTED_SLIDES} titles are #{EXPECTED_COLOR} ({color_score:.2f} pts)")
            total_score += color_score
        else:
            print(f"FAIL: Component 3 -- 0/{EXPECTED_SLIDES} titles are #{EXPECTED_COLOR}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
