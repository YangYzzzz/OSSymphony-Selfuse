"""
Initial Setup: Financial presentation with FY2024 references
Task ID: impstruct_017
Domain: libreoffice_impress

Creates a 10-slide financial update presentation with ~15 occurrences of 'FY2024'
in titles, body text, and speaker notes.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
OUTPUT = f'{WORKDIR}/financial_update.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_shape(shape, text, font_size=18, bold=False, color=None, alignment=None):
    """Helper to set text on a shape with formatting."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    for run in p.runs:
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color


def add_paragraph(text_frame, text, font_size=14, bold=False, color=None, level=0):
    """Add a paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    for run in p.runs:
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
    return p


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    # FY2024 occurrence #1 (title)
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Nextera Solutions — FY2024 Financial Update"
    slide1.placeholders[1].text = "Quarterly Business Review | December 2024"

    # --- Slide 2: Executive Summary ---
    # FY2024 occurrence #2 (title), #3 (body)
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "FY2024 Executive Summary"
    body2 = slide2.placeholders[1].text_frame
    body2.paragraphs[0].text = "Total revenue reached $142.8M, exceeding targets by 8%"
    add_paragraph(body2, "Operating margin improved to 23.4%, up from 19.1% last year")
    add_paragraph(body2, "FY2024 headcount grew by 340 employees across 6 regions")
    add_paragraph(body2, "Customer retention rate held steady at 94.2%")
    slide2.notes_slide.notes_text_frame.text = "Emphasize the strong performance in FY2024 Q3 and Q4"

    # --- Slide 3: Revenue Breakdown ---
    # FY2024 occurrence #4 (title)
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Revenue Breakdown — FY2024"
    body3 = slide3.placeholders[1].text_frame
    body3.paragraphs[0].text = "Enterprise Solutions: $68.5M (48%)"
    add_paragraph(body3, "Cloud Services: $41.2M (29%)")
    add_paragraph(body3, "Professional Services: $22.4M (16%)")
    add_paragraph(body3, "Licensing & Support: $10.7M (7%)")

    # --- Slide 4: Regional Performance ---
    # FY2024 occurrence #5 (body), #6 (notes)
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Regional Performance"
    body4 = slide4.placeholders[1].text_frame
    body4.paragraphs[0].text = "North America led FY2024 revenue growth at 12.3% YoY"
    add_paragraph(body4, "EMEA region delivered $38.9M with strong expansion in Germany")
    add_paragraph(body4, "APAC contributed $24.1M, with Japan as the top market")
    add_paragraph(body4, "Latin America emerging market pilot launched in Brazil")
    slide4.notes_slide.notes_text_frame.text = "FY2024 regional breakdown shows North America still dominant but APAC gaining share"

    # --- Slide 5: Product Highlights ---
    # FY2024 occurrence #7 (body)
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Product Highlights"
    body5 = slide5.placeholders[1].text_frame
    body5.paragraphs[0].text = "Platform v4.2 released in March with AI-driven analytics"
    add_paragraph(body5, "Mobile app downloads surpassed 1.2M in FY2024")
    add_paragraph(body5, "API integrations expanded to 85+ third-party connectors")
    add_paragraph(body5, "NPS score improved to 72, up from 64")

    # --- Slide 6: Operational Efficiency ---
    # FY2024 occurrence #8 (title), #9 (notes)
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "FY2024 Operational Efficiency Metrics"
    body6 = slide6.placeholders[1].text_frame
    body6.paragraphs[0].text = "Infrastructure costs reduced by 15% through cloud migration"
    add_paragraph(body6, "Average ticket resolution time: 4.2 hours (down from 6.8)")
    add_paragraph(body6, "Employee productivity index: 1.34 (target was 1.25)")
    add_paragraph(body6, "System uptime maintained at 99.97%")
    slide6.notes_slide.notes_text_frame.text = "Compare FY2024 OpEx improvements to industry benchmarks"

    # --- Slide 7: Talent & Culture ---
    # FY2024 occurrence #10 (body)
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Talent & Culture"
    body7 = slide7.placeholders[1].text_frame
    body7.paragraphs[0].text = "Total headcount reached 2,840 employees"
    add_paragraph(body7, "Voluntary turnover rate decreased to 8.3% in FY2024")
    add_paragraph(body7, "Diversity hiring improved: 47% of new hires from underrepresented groups")
    add_paragraph(body7, "Launched leadership development program for 120 managers")

    # --- Slide 8: Risk Assessment ---
    # FY2024 occurrence #11 (body), #12 (notes)
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Risk Assessment"
    body8 = slide8.placeholders[1].text_frame
    body8.paragraphs[0].text = "Currency exposure managed within 2% variance band"
    add_paragraph(body8, "Cybersecurity investment increased 30% in FY2024")
    add_paragraph(body8, "Supply chain diversification reduced single-vendor dependency to 12%")
    add_paragraph(body8, "Regulatory compliance maintained across all operating jurisdictions")
    slide8.notes_slide.notes_text_frame.text = "FY2024 risk profile improved overall; key residual risk is FX exposure in emerging markets"

    # --- Slide 9: Strategic Outlook ---
    # FY2024 occurrence #13 (body), #14 (body)
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Strategic Outlook"
    body9 = slide9.placeholders[1].text_frame
    body9.paragraphs[0].text = "Building on FY2024 momentum to target $165M revenue"
    add_paragraph(body9, "Three new product lines planned for launch in H1")
    add_paragraph(body9, "Expanding APAC presence with offices in Singapore and Seoul")
    add_paragraph(body9, "FY2024 lessons learned integrated into updated strategy framework")

    # --- Slide 10: Thank You ---
    # FY2024 occurrence #15 (notes)
    slide10 = prs.slides.add_slide(prs.slide_layouts[0])
    slide10.shapes.title.text = "Thank You"
    slide10.placeholders[1].text = "Questions & Discussion"
    slide10.notes_slide.notes_text_frame.text = "Wrap up with key FY2024 takeaways and open the floor for questions"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Count FY2024 occurrences for verification
    count = 0
    prs2 = Presentation(OUTPUT)
    for slide in prs2.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        count += run.text.count('FY2024')
        try:
            notes = slide.notes_slide.notes_text_frame
            for para in notes.paragraphs:
                for run in para.runs:
                    count += run.text.count('FY2024')
        except Exception:
            pass
    print(f'FY2024 occurrences found: {count}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
