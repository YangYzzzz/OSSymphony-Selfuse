"""
Reward Script: Find & Replace FY2024 -> FY2025 throughout presentation
Task ID: impstruct_017
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): No 'FY2024' remains in shape text across all slides
  Component 2 (0.3): 'FY2025' appears in shape text (>= 10 occurrences expected)
  Component 3 (0.3): Notes also updated — no 'FY2024' in notes AND 'FY2025' present in notes
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impstruct_017'


def persist_app_state(domain: str):
    """Save any unsaved LibreOffice edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify that all occurrences of 'FY2024' have been replaced with 'FY2025'
    in both slide shapes and notes.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Gather all text from shapes and notes
    shape_fy2024_count = 0
    shape_fy2025_count = 0
    notes_fy2024_count = 0
    notes_fy2025_count = 0

    for i, slide in enumerate(prs.slides):
        # Check shape text (including grouped shapes)
        def extract_text_shapes(shape):
            shapes = []
            if hasattr(shape, 'text') and hasattr(shape, 'text_frame'):
                shapes.append(shape)
            if hasattr(shape, 'shapes'):
                for sub in shape.shapes:
                    shapes.extend(extract_text_shapes(sub))
            return shapes

        for shape in slide.shapes:
            for text_shape in extract_text_shapes(shape):
                for para in text_shape.text_frame.paragraphs:
                    text = para.text
                    shape_fy2024_count += text.count('FY2024')
                    shape_fy2025_count += text.count('FY2025')

        # Check notes
        try:
            notes_text = slide.notes_slide.notes_text_frame.text
            notes_fy2024_count += notes_text.count('FY2024')
            notes_fy2025_count += notes_text.count('FY2025')
        except Exception:
            pass

    print(f"Shape FY2024 count: {shape_fy2024_count}")
    print(f"Shape FY2025 count: {shape_fy2025_count}")
    print(f"Notes FY2024 count: {notes_fy2024_count}")
    print(f"Notes FY2025 count: {notes_fy2025_count}")

    # Component 1: No 'FY2024' remains in shape text (0.4 points)
    # Initial has 12 shape occurrences; golden has 0
    try:
        if shape_fy2024_count == 0:
            print(f"PASS: Component 1 -- No FY2024 in shapes (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 -- Found {shape_fy2024_count} FY2024 in shapes, expected 0")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: 'FY2025' appears in shape text with sufficient occurrences (0.3 points)
    # Golden has 12 shape occurrences; initial has 0
    try:
        if shape_fy2025_count >= 10:
            print(f"PASS: Component 2 -- Found {shape_fy2025_count} FY2025 in shapes (0.3 pts)")
            total_score += 0.3
        elif shape_fy2025_count > 0:
            # Partial credit proportional to how many were replaced
            partial = 0.3 * (shape_fy2025_count / 12.0)
            print(f"PARTIAL: Component 2 -- Found {shape_fy2025_count}/12 FY2025 in shapes ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 -- No FY2025 found in shapes, expected ~12")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Notes updated — no FY2024 in notes AND FY2025 present in notes (0.3 points)
    # Initial has 4 FY2024 in notes, 0 FY2025; golden has 0 FY2024, 4 FY2025
    try:
        no_old_in_notes = (notes_fy2024_count == 0)
        new_in_notes = (notes_fy2025_count >= 3)

        if no_old_in_notes and new_in_notes:
            print(f"PASS: Component 3 -- Notes updated: 0 FY2024, {notes_fy2025_count} FY2025 (0.3 pts)")
            total_score += 0.3
        elif no_old_in_notes or new_in_notes:
            # Half credit if one condition met
            print(f"PARTIAL: Component 3 -- Notes partially updated: FY2024={notes_fy2024_count}, FY2025={notes_fy2025_count} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 3 -- Notes not updated: FY2024={notes_fy2024_count}, FY2025={notes_fy2025_count}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist unsaved GUI state before checking
persist_app_state("libreoffice_impress")

# Default: test against canonical artifact path
file_path = f'{WORKDIR}/financial_update.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
