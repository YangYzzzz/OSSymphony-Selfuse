"""
Reward Script: Add Dissolve transitions to all slides and insert Thank You slide
Task ID: impress_wf_005
Domain: libreoffice_impress
Scoring:
  Component 1: Slide count == 6           (0.15 pts)
  Component 2: All 6 slides have Dissolve (0.30 pts)
  Component 3: All 6 slides have 2s dur   (0.15 pts)
  Component 4: Slide 6 bg #1B5E20        (0.15 pts)
  Component 5: Slide 6 text properties    (0.25 pts)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_005'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: file exists and is loadable
    try:
        from pptx import Presentation
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Slide count is 6 (0.15 points)
    # Initial has 5 slides, golden should have 6 (5 original + 1 Thank You)
    try:
        if num_slides == 6:
            print(f"PASS: Component 1 — Slide count is 6 (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — Expected 6 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: All slides have Dissolve transition (0.30 points)
    # Initial has NO transitions on any slide; golden has Dissolve on all 6
    try:
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        dissolve_count = 0
        with zipfile.ZipFile(file_path, 'r') as zf:
            for i in range(1, num_slides + 1):
                try:
                    with zf.open(f'ppt/slides/slide{i}.xml') as f:
                        root = ET.parse(f).getroot()
                        tr = root.find('.//p:transition', ns)
                        if tr is not None and tr.find('.//p:dissolve', ns) is not None:
                            dissolve_count += 1
                except KeyError:
                    pass

        if num_slides >= 6 and dissolve_count == num_slides:
            print(f"PASS: Component 2 — All {num_slides} slides have Dissolve transition (0.30 pts)")
            total_score += 0.30
        elif dissolve_count > 0:
            # Partial credit: proportional to how many slides have dissolve
            partial = 0.30 * (dissolve_count / max(num_slides, 6))
            print(f"PARTIAL: Component 2 — {dissolve_count}/{num_slides} slides have Dissolve ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No slides have Dissolve transition")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: All slides have 2-second transition duration (0.15 points)
    # Initial has no transitions; golden has dur="2000" on all slides
    try:
        dur_correct_count = 0
        with zipfile.ZipFile(file_path, 'r') as zf:
            for i in range(1, num_slides + 1):
                try:
                    with zf.open(f'ppt/slides/slide{i}.xml') as f:
                        root = ET.parse(f).getroot()
                        tr = root.find('.//p:transition', ns)
                        if tr is not None:
                            dur = tr.get('dur')
                            # 2 seconds = 2000 ms
                            if dur is not None and int(dur) == 2000:
                                dur_correct_count += 1
                except (KeyError, ValueError):
                    pass

        if num_slides >= 6 and dur_correct_count == num_slides:
            print(f"PASS: Component 3 — All {num_slides} slides have 2s duration (0.15 pts)")
            total_score += 0.15
        elif dur_correct_count > 0:
            partial = 0.15 * (dur_correct_count / max(num_slides, 6))
            print(f"PARTIAL: Component 3 — {dur_correct_count}/{num_slides} slides have 2s duration ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No slides have 2-second duration")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Last slide (slide 6) has dark green background #1B5E20 (0.15 points)
    # This slide doesn't exist in initial_env, so this check inherently fails on initial
    try:
        if num_slides >= 6:
            last_slide = prs.slides[num_slides - 1]
            fill = last_slide.background.fill
            bg_color = None
            if fill.type == 1:  # SOLID fill
                bg_color = str(fill.fore_color.rgb)

            if bg_color is not None and bg_color.upper() == '1B5E20':
                print(f"PASS: Component 4 — Last slide background is #1B5E20 (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 4 — Expected bg #1B5E20, found type={fill.type}, color={bg_color}")
        else:
            print(f"FAIL: Component 4 — Not enough slides ({num_slides}) to check slide 6 background")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Last slide text "Thank You!" centered, 54pt bold white (0.25 points)
    # Sub-checks: text content (0.10), bold+size (0.08), white color+centered (0.07)
    try:
        if num_slides >= 6:
            last_slide = prs.slides[num_slides - 1]
            found_thank_you = False
            text_correct = False
            format_correct = False
            color_align_correct = False

            for shape in last_slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if 'thank you' in para.text.strip().lower():
                            found_thank_you = True
                            text_correct = para.text.strip() == 'Thank You!'

                            # Check alignment (CENTER)
                            alignment = para.alignment
                            is_centered = (alignment is not None and
                                           alignment == PP_ALIGN.CENTER)

                            for run in para.runs:
                                if 'thank you' in run.text.lower():
                                    # Check bold
                                    is_bold = run.font.bold is True

                                    # Check size (54pt = 685800 EMU)
                                    is_54pt = (run.font.size is not None and
                                               abs(run.font.size - 685800) < 1000)

                                    # Check white color
                                    is_white = False
                                    try:
                                        if run.font.color.type is not None:
                                            rgb = str(run.font.color.rgb).upper()
                                            is_white = rgb == 'FFFFFF'
                                    except Exception:
                                        pass

                                    format_correct = is_bold and is_54pt
                                    color_align_correct = is_white and is_centered
                                    break
                            break

            if text_correct:
                print(f"PASS: Component 5a — Text is 'Thank You!' (0.10 pts)")
                total_score += 0.10
            elif found_thank_you:
                print(f"PARTIAL: Component 5a — Found thank you text but not exact match (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 5a — 'Thank You!' text not found on last slide")

            if format_correct:
                print(f"PASS: Component 5b — Bold and 54pt (0.08 pts)")
                total_score += 0.08
            else:
                print(f"FAIL: Component 5b — Expected bold=True and size=54pt")

            if color_align_correct:
                print(f"PASS: Component 5c — White color and centered (0.07 pts)")
                total_score += 0.07
            else:
                print(f"FAIL: Component 5c — Expected white (#FFFFFF) and centered alignment")
        else:
            print(f"FAIL: Component 5 — Not enough slides ({num_slides}) to check last slide text")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
