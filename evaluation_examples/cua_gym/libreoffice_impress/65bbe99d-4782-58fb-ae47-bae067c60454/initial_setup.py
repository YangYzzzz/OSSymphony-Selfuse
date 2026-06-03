"""
Initial Setup: Create Meeting_Notes.pptx with 5 slides, no transitions.
Task ID: impress_wf_005
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_005'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, alignment=PP_ALIGN.LEFT, color=None):
    """Helper to add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tf


def add_bullet_list(text_frame, items, font_size=14, color=None):
    """Add bullet items to an existing text frame."""
    for item in items:
        p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        run = p.runs[0]
        run.font.size = Pt(font_size)
        if color:
            run.font.color.rgb = color


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q2 Strategy Meeting Notes"
    slide1.placeholders[1].text = "Marketing & Product Division\nApril 10, 2025"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    add_text_box(slide2, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
                 "Meeting Agenda", font_size=28, bold=True,
                 color=RGBColor(0x1F, 0x49, 0x7D))
    tf2 = add_text_box(slide2, Inches(0.8), Inches(1.5), Inches(8), Inches(4.5),
                       "1. Review Q1 campaign performance metrics", font_size=16)
    add_bullet_list(tf2, [
        "2. Discuss product launch timeline for Project Aurora",
        "3. Budget reallocation for digital advertising channels",
        "4. Cross-team collaboration updates with Engineering",
        "5. Customer feedback analysis from March surveys",
        "6. Open floor for questions and suggestions",
    ], font_size=16)

    # --- Slide 3: Key Discussion Points ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    add_text_box(slide3, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
                 "Key Discussion Points", font_size=28, bold=True,
                 color=RGBColor(0x1F, 0x49, 0x7D))
    tf3 = add_text_box(slide3, Inches(0.8), Inches(1.5), Inches(8), Inches(5),
                       "Q1 Performance: Revenue up 12% YoY, exceeding $2.4M target",
                       font_size=14)
    add_bullet_list(tf3, [
        "Project Aurora: Beta testing scheduled for June 15, 2025",
        "Digital ad spend: Reallocating $150K from print to social media",
        "Engineering sync: API integration on track for July release",
        "Customer NPS improved from 42 to 58 after UX redesign",
        "Churn rate decreased by 3.2% in enterprise segment",
        "New partnership with TechVista Solutions confirmed for Q3",
    ], font_size=14)

    # --- Slide 4: Action Items ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    add_text_box(slide4, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
                 "Action Items", font_size=28, bold=True,
                 color=RGBColor(0x1F, 0x49, 0x7D))

    # Add a table for action items
    rows, cols = 6, 3
    table_shape = slide4.shapes.add_table(rows, cols,
                                          Inches(0.8), Inches(1.5),
                                          Inches(8), Inches(3.5))
    table = table_shape.table
    table.columns[0].width = Inches(4)
    table.columns[1].width = Inches(2)
    table.columns[2].width = Inches(2)

    headers = ["Action Item", "Owner", "Deadline"]
    data = [
        ["Finalize Q2 marketing budget proposal", "Sarah Chen", "Apr 18, 2025"],
        ["Schedule beta tester recruitment drive", "Marcus Johnson", "Apr 25, 2025"],
        ["Prepare social media campaign assets", "Priya Patel", "May 2, 2025"],
        ["Draft TechVista partnership agreement", "David Kim", "Apr 22, 2025"],
        ["Update customer onboarding flow mockups", "Elena Rodriguez", "May 9, 2025"],
    ]

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(12)

    # --- Slide 5: Next Steps ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    add_text_box(slide5, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
                 "Next Steps & Timeline", font_size=28, bold=True,
                 color=RGBColor(0x1F, 0x49, 0x7D))
    tf5 = add_text_box(slide5, Inches(0.8), Inches(1.5), Inches(8), Inches(5),
                       "Week of Apr 14: Budget review with Finance team",
                       font_size=14)
    add_bullet_list(tf5, [
        "Week of Apr 21: Beta tester outreach begins",
        "Week of Apr 28: Creative assets due for social campaigns",
        "Week of May 5: Engineering API milestone check-in",
        "Week of May 12: Mid-quarter progress review",
        "Next full team meeting: May 15, 2025 at 2:00 PM",
    ], font_size=14)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
