"""
Initial Setup: Web Analytics presentation with 6 slides, slide 4 titled 'Monthly Website Traffic' with no chart.
Task ID: impress_tct_038
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_038'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_text(slide, text):
    """Set the title placeholder text."""
    if slide.shapes.title:
        slide.shapes.title.text = text


def add_body_text(slide, paragraphs):
    """Add text to the first non-title placeholder (body/content area)."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0:  # skip title
            tf = ph.text_frame
            tf.word_wrap = True
            for i, (txt, level) in enumerate(paragraphs):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = txt
                p.level = level
                for run in p.runs:
                    run.font.size = Pt(16)
            break


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Web Analytics Report"
    slide1.placeholders[1].text = "Q1-Q2 2025 Performance Review"

    # --- Slide 2: Traffic Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide2, "Traffic Overview")
    add_body_text(slide2, [
        ("Total sessions increased by 34% compared to last quarter", 0),
        ("Mobile traffic now accounts for 62% of all visits", 0),
        ("Average session duration: 3 minutes 42 seconds", 0),
        ("Bounce rate decreased from 45% to 38%", 0),
        ("Top referral sources: Google Search (41%), Direct (28%), Social Media (18%)", 0),
    ])

    # --- Slide 3: User Demographics ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide3, "User Demographics")
    add_body_text(slide3, [
        ("Age distribution of website visitors:", 0),
        ("18-24: 22% of total users", 1),
        ("25-34: 35% of total users", 1),
        ("35-44: 24% of total users", 1),
        ("45-54: 12% of total users", 1),
        ("55+: 7% of total users", 1),
        ("Geographic breakdown: North America 48%, Europe 31%, Asia-Pacific 15%, Other 6%", 0),
    ])

    # --- Slide 4: Monthly Website Traffic (NO CHART - just title) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Monthly Website Traffic"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # --- Slide 5: Conversion Rates ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide5, "Conversion Rates")
    add_body_text(slide5, [
        ("Overall conversion rate: 3.8% (up from 2.9%)", 0),
        ("Landing page optimization improved sign-ups by 22%", 0),
        ("Cart abandonment rate reduced to 61% from 73%", 0),
        ("Email campaign CTR: 4.2% (industry avg: 2.6%)", 0),
        ("A/B test results show new CTA button increased clicks by 17%", 0),
    ])

    # --- Slide 6: Summary & Next Steps ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide6, "Summary & Next Steps")
    add_body_text(slide6, [
        ("Key achievements this quarter:", 0),
        ("Record-breaking traffic in May with 11,000 unique visitors", 1),
        ("Mobile experience overhaul contributed to lower bounce rates", 1),
        ("SEO improvements drove 34% organic traffic growth", 1),
        ("Next steps:", 0),
        ("Implement personalized content recommendations", 1),
        ("Expand social media advertising budget by 20%", 1),
        ("Launch redesigned mobile checkout flow in Q3", 1),
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
