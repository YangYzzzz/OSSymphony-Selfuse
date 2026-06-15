"""
Reward Script: Insert a line chart on slide 4 with monthly website traffic data
Task ID: impress_tct_038
Domain: libreoffice_impress
Scoring:
  Component 1: Chart exists on slide 4 (0.2 pts)
  Component 2: Chart type is LINE (0.2 pts)
  Component 3: Categories are Jan-Jun (0.2 pts)
  Component 4: Data values match expected [5000, 7200, 6800, 9100, 8500, 11000] (0.3 pts)
  Component 5: Exactly one data series (0.1 pts)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_038'


def persist_app_state(domain):
    """Send Ctrl+S to save any unsaved GUI edits."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.chart import XL_CHART_TYPE
    except ImportError as e:
        print(f"CRITICAL: Missing dependency: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 4 slides
    if len(prs.slides) < 4:
        print(f"PRECONDITION FAIL: Expected at least 4 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide4 = prs.slides[3]  # 0-indexed

    # Find chart shapes on slide 4
    chart_shapes = [s for s in slide4.shapes if s.has_chart]

    # Component 1: Chart exists on slide 4 (0.2 points)
    try:
        if len(chart_shapes) > 0:
            print(f"PASS: Component 1 -- Chart found on slide 4 ({len(chart_shapes)} chart(s)) (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 1 -- No chart found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    if len(chart_shapes) == 0:
        # No chart means all remaining checks fail
        print("FAIL: Components 2-5 skipped (no chart)")
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    chart = chart_shapes[0].chart

    # Component 2: Chart type is LINE (0.2 points)
    try:
        chart_type = chart.chart_type
        # LINE = 4, LINE_MARKERS = 8, LINE_MARKERS_STACKED = 9, LINE_STACKED = 5
        line_types = {
            XL_CHART_TYPE.LINE,
            XL_CHART_TYPE.LINE_MARKERS,
            XL_CHART_TYPE.LINE_STACKED,
            XL_CHART_TYPE.LINE_MARKERS_STACKED,
        }
        if chart_type in line_types:
            print(f"PASS: Component 2 -- Chart type is LINE variant ({chart_type}) (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 2 -- Expected LINE type, found {chart_type}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Categories match Jan-Jun (0.2 points)
    expected_categories = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']
    try:
        actual_categories = [str(c) for c in chart.plots[0].categories]
        if actual_categories == expected_categories:
            print(f"PASS: Component 3 -- Categories match {expected_categories} (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 3 -- Expected {expected_categories}, found {actual_categories}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Data values match expected (0.3 points)
    expected_values = [5000.0, 7200.0, 6800.0, 9100.0, 8500.0, 11000.0]
    try:
        series_list = list(chart.series)
        if len(series_list) >= 1:
            actual_values = list(series_list[0].values)
            # Compare with tolerance for floating point
            if len(actual_values) == len(expected_values):
                all_match = all(
                    abs(a - e) < 1.0
                    for a, e in zip(actual_values, expected_values)
                )
                if all_match:
                    print(f"PASS: Component 4 -- Data values match {expected_values} (0.3 pts)")
                    total_score += 0.3
                else:
                    print(f"FAIL: Component 4 -- Values mismatch. Expected {expected_values}, found {actual_values}")
            else:
                print(f"FAIL: Component 4 -- Expected {len(expected_values)} data points, found {len(actual_values)}")
        else:
            print(f"FAIL: Component 4 -- No series found in chart")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Exactly one data series (0.1 points)
    try:
        num_series = len(list(chart.series))
        if num_series == 1:
            print(f"PASS: Component 5 -- Exactly 1 data series (0.1 pts)")
            total_score += 0.1
        else:
            print(f"FAIL: Component 5 -- Expected 1 series, found {num_series}")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist app state before verification (LibreOffice may have unsaved edits)
persist_app_state("libreoffice_impress")

# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
