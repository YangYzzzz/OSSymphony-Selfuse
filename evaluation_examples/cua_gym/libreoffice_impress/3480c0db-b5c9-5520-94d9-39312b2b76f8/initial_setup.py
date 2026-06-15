"""
Initial Setup: Business Plan presentation with empty SWOT Analysis slide
Task ID: impress_stu_058
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_058'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_body(slide, title_text, body_lines):
    """Helper to set title and body content on a standard slide."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    # Find the body placeholder (index 1 typically)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.text = body_lines[0]
            for line in body_lines[1:]:
                p = tf.add_paragraph()
                p.text = line
                p.level = 0
            break


def create_initial():
    prs = Presentation()

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Campus Bookstore Business Plan"
    slide1.placeholders[1].text = "Prepared by the Student Business Association\nSpring 2025"

    # ---- Slide 2: Executive Summary ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide2, "Executive Summary", [
        "The Campus Bookstore serves over 12,000 students annually",
        "Revenue reached $2.4M in the last fiscal year",
        "Expanding into digital course materials and branded merchandise",
        "Targeting 15% revenue growth through online platform integration",
        "Strategic partnerships with 3 major publishers secured",
    ])

    # ---- Slide 3: Market Analysis ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide3, "Market Analysis", [
        "Total addressable market: 15,200 enrolled students + 2,800 faculty",
        "Current market penetration: 68% for textbooks, 42% for supplies",
        "Key competitor: Online retailers capturing 30% of textbook sales",
        "Growing demand for rental programs and digital access codes",
        "Survey data: 74% of students prefer on-campus pickup options",
    ])

    # ---- Slide 4: SWOT Analysis (EMPTY — task target) ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add only the title as a text box at the top
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "SWOT Analysis"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    # Body is intentionally EMPTY — the agent must create the SWOT grid

    # ---- Slide 5: Marketing Strategy ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide5, "Marketing Strategy", [
        "Launch loyalty rewards program for repeat customers",
        "Social media campaigns targeting freshmen during orientation week",
        "Partner with student organizations for co-branded events",
        "Introduce price-match guarantee for top 50 textbook titles",
        "Email marketing to faculty for bulk course material orders",
    ])

    # ---- Slide 6: Financial Projections ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide6, "Financial Projections", [
        "Year 1 projected revenue: $2.76M (15% growth)",
        "Operating margin target: 18% (up from 14%)",
        "Capital expenditure: $180K for online platform development",
        "Break-even on digital investment expected within 14 months",
        "Three-year revenue target: $3.5M with diversified income streams",
    ])

    # ---- Slide 7: Operations Plan ----
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide7, "Operations Plan", [
        "Extend store hours during midterms and finals (7 AM - 11 PM)",
        "Hire 8 additional part-time student employees for peak seasons",
        "Implement inventory management system (NetSuite integration)",
        "Establish textbook buyback kiosk in student union building",
        "Monthly staff training on customer service best practices",
    ])

    # ---- Slide 8: Conclusion ----
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide8, "Conclusion & Next Steps", [
        "The campus bookstore is well-positioned for sustainable growth",
        "Digital transformation is critical to compete with online retailers",
        "Investment in student experience drives long-term loyalty",
        "Board approval requested for $180K platform development budget",
        "Implementation timeline: Q3 2025 launch of online ordering system",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
