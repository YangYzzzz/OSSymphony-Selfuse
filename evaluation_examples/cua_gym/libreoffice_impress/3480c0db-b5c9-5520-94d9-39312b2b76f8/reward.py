"""
Reward Script: SWOT Analysis diagram on slide 4
Task ID: impress_stu_058
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): 4 rounded rectangle shapes on slide 4
  Component 2 (0.25): Correct fill colors (green, orange, blue, red)
  Component 3 (0.25): Correct titles (Strengths/Weaknesses/Opportunities/Threats) bold, ~20pt, white
  Component 4 (0.25): Each rectangle has 3 bullet items in white, ~14pt
"""

import os
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_058'

# Expected SWOT colors (uppercase hex without #)
EXPECTED_COLORS = {
    'Strengths': '27AE60',
    'Weaknesses': 'E67E22',
    'Opportunities': '2980B9',
    'Threats': 'C0392B',
}


def get_auto_shapes_on_slide(slide):
    """Return all AUTO_SHAPE shapes on a slide (type 1)."""
    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]


def get_shape_fill_color(shape):
    """Get the solid fill color of a shape as uppercase hex string, or None."""
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb).upper()
    except Exception:
        pass
    return None


def get_shape_title_and_bullets(shape):
    """Extract the title (first non-empty paragraph) and bullet texts from a shape.
    Returns (title_text, title_run_info, bullet_texts, bullet_run_infos).
    title_run_info: dict with bold, size, color_rgb
    bullet_run_infos: list of dicts
    """
    if not shape.has_text_frame:
        return None, None, [], []

    title_text = None
    title_info = None
    bullet_texts = []
    bullet_infos = []

    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # Gather run info from first non-empty run
        run_info = None
        for run in para.runs:
            if (run.text or "").strip():
                info = {
                    'bold': run.font.bold if run.font.bold is not None else False,
                    'size': run.font.size,
                    'color_rgb': None,
                }
                try:
                    if run.font.color.type is not None:
                        info['color_rgb'] = str(run.font.color.rgb).upper()
                except Exception:
                    pass
                run_info = info
                break

        if title_text is None:
            # First non-empty paragraph is the title
            title_text = text
            title_info = run_info
        else:
            # Subsequent paragraphs are bullets
            # Strip leading bullet character if present
            clean = text.lstrip('•●○■◆- ').strip()
            bullet_texts.append(clean)
            bullet_infos.append(run_info)

    return title_text, title_info, bullet_texts, bullet_infos


def verify_task(file_path):
    """
    Verify SWOT analysis diagram on slide 4.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Check we have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)

    # Get auto shapes (rounded rectangles are AUTO_SHAPE type)
    auto_shapes = get_auto_shapes_on_slide(slide)

    # Component 1: 4 rounded rectangle auto shapes on slide 4 (0.25 points)
    try:
        num_auto = len(auto_shapes)
        if num_auto >= 4:
            print(f"PASS: Component 1 — Found {num_auto} auto shapes on slide 4 (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — Expected 4 auto shapes, found {num_auto}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Build a mapping: title -> (shape, fill_color, title_info, bullets, bullet_infos)
    shape_map = {}
    for shape in auto_shapes:
        fill_color = get_shape_fill_color(shape)
        title_text, title_info, bullet_texts, bullet_infos = get_shape_title_and_bullets(shape)
        if title_text:
            shape_map[title_text] = {
                'fill_color': fill_color,
                'title_info': title_info,
                'bullet_texts': bullet_texts,
                'bullet_infos': bullet_infos,
            }

    # Component 2: Correct fill colors for all 4 SWOT sections (0.25 points)
    try:
        color_matches = 0
        for title, expected_hex in EXPECTED_COLORS.items():
            if title in shape_map:
                actual_color = shape_map[title].get('fill_color')
                if actual_color and actual_color.upper() == expected_hex.upper():
                    print(f"  PASS: {title} fill color = {actual_color}")
                    color_matches += 1
                else:
                    print(f"  FAIL: {title} expected fill {expected_hex}, got {actual_color}")
            else:
                print(f"  FAIL: No shape titled '{title}' found")

        if color_matches == 4:
            print(f"PASS: Component 2 — All 4 SWOT fill colors correct (0.25 pts)")
            total_score += 0.25
        elif color_matches >= 2:
            partial = round(0.25 * color_matches / 4, 2)
            print(f"PARTIAL: Component 2 — {color_matches}/4 colors correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — Only {color_matches}/4 fill colors correct")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Titles are bold, ~20pt, white text (0.25 points)
    try:
        title_matches = 0
        for title in EXPECTED_COLORS:
            if title not in shape_map:
                print(f"  FAIL: Title '{title}' not found")
                continue

            info = shape_map[title].get('title_info')
            if info is None:
                print(f"  FAIL: No run info for title '{title}'")
                continue

            checks_ok = True
            # Bold check
            if not info['bold']:
                print(f"  FAIL: '{title}' title not bold")
                checks_ok = False

            # Size check: ~20pt = 254000 EMU (allow 16pt-24pt range)
            if info['size'] is not None:
                size_pt = info['size'] / 12700  # EMU to pt
                if 16 <= size_pt <= 28:
                    pass  # acceptable
                else:
                    print(f"  FAIL: '{title}' title size {size_pt}pt, expected ~20pt")
                    checks_ok = False

            # White text check
            if info['color_rgb'] and info['color_rgb'] == 'FFFFFF':
                pass  # white
            else:
                print(f"  FAIL: '{title}' title color {info['color_rgb']}, expected FFFFFF")
                checks_ok = False

            if checks_ok:
                title_matches += 1
                print(f"  PASS: '{title}' title is bold, ~20pt, white")

        if title_matches == 4:
            print(f"PASS: Component 3 — All 4 titles correct (0.25 pts)")
            total_score += 0.25
        elif title_matches >= 2:
            partial = round(0.25 * title_matches / 4, 2)
            print(f"PARTIAL: Component 3 — {title_matches}/4 titles correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — Only {title_matches}/4 titles correct")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Each rectangle has 3 bullet items, white, ~14pt (0.25 points)
    try:
        bullet_matches = 0
        for title in EXPECTED_COLORS:
            if title not in shape_map:
                print(f"  FAIL: '{title}' not found for bullet check")
                continue

            bullets = shape_map[title]['bullet_texts']
            b_infos = shape_map[title]['bullet_infos']

            if len(bullets) < 3:
                print(f"  FAIL: '{title}' has {len(bullets)} bullets, expected 3")
                continue

            # Check that bullets have white text and ~14pt
            all_ok = True
            for idx, binfo in enumerate(b_infos[:3]):
                if binfo is None:
                    continue
                # White check
                if binfo['color_rgb'] and binfo['color_rgb'] != 'FFFFFF':
                    print(f"  FAIL: '{title}' bullet {idx+1} color {binfo['color_rgb']}, expected FFFFFF")
                    all_ok = False
                # Size check: ~14pt = 177800 EMU (allow 12-18pt)
                if binfo['size'] is not None:
                    size_pt = binfo['size'] / 12700
                    if not (10 <= size_pt <= 20):
                        print(f"  FAIL: '{title}' bullet {idx+1} size {size_pt}pt, expected ~14pt")
                        all_ok = False

            if all_ok:
                bullet_matches += 1
                print(f"  PASS: '{title}' has 3+ bullets, white, ~14pt")

        if bullet_matches == 4:
            print(f"PASS: Component 4 — All 4 sections have correct bullets (0.25 pts)")
            total_score += 0.25
        elif bullet_matches >= 2:
            partial = round(0.25 * bullet_matches / 4, 2)
            print(f"PARTIAL: Component 4 — {bullet_matches}/4 sections correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — Only {bullet_matches}/4 sections have correct bullets")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
