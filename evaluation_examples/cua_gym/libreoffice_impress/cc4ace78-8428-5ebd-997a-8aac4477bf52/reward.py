"""
FINAL REWARD SCRIPT - SUCCESS
Task: I'm polishing a 50-slide deck and need slide 38 to stand out. In LibreOffice Impress, how can I swap the plain background on that one slide for the built-in "Gradient Blue 2" background without changing any of the other slides?
Generated: 2025-09-10 14:51:07
Status: success
Model: azure-o3
Total Steps: 3
"""

from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.dml.color import RGBColor
import os
import traceback

"""
Reward Script
Task: Verify that in the provided 50-slide deck, ONLY slide 38 has had its
plain background swapped for the built-in "Gradient Blue 2" background, while
all other slides remain unchanged (i.e., they do NOT use a gradient fill).

Scoring (progressive):
• 0.4 pts – Presentation contains exactly 50 slides (ensures correct file)
• 0.4 pts – Slide 38 background is a gradient AND its two stop colours match
               LibreOffice’s built-in "Gradient Blue 2" (4472C4 ➜ D9E1F2)
• 0.2 pts – No other slide (except #38) uses a gradient background
Total = 1.0 when all checks pass.

The script prints detailed diagnostics and finally outputs
"REWARD: X.X" (float between 0.0-1.0).
"""

FILE_PATH = (
    "/home/user/"
    "im_polishing_a_50_slide_deck_and_need_slide_38_to_stand_out_in_libreoffice_impress_how_can_i_swap_th_golden.pptx"
)

EXPECTED_SLIDE_COUNT = 50          # deck size should remain unchanged
TARGET_SLIDE_INDEX = 37            # zero-based index → slide 38
GRADIENT_BLUE2_COLORS = {
    RGBColor(0x44, 0x72, 0xC4),    # dark blue left/top
    RGBColor(0xD9, 0xE1, 0xF2),    # light blue right/bottom
}


def verify_task(file_path: str) -> float:
    """Return a float score ∈ [0,1] based on task completion."""
    total_score = 0.0
    MAX_SCORE = 1.0

    # ---------- 1. File existence & load (no points) ----------
    if not os.path.exists(file_path):
        print("✗ File not found:", file_path)
        return 0.0  # cannot continue verification

    try:
        prs = Presentation(file_path)
    except Exception as exc:
        print("✗ Could not open PPTX:", exc)
        traceback.print_exc()
        return 0.0

    # ---------- 2. Slide count check (0.4 pts) ----------
    slide_count = len(prs.slides)
    print(f"Slide count → expected {EXPECTED_SLIDE_COUNT}, found {slide_count}")
    if slide_count == EXPECTED_SLIDE_COUNT:
        print("✓ Slide count matches (0.4 pts)")
        total_score += 0.4
    else:
        print("✗ Slide count mismatch (0 pts)")

    # ---------- 3. Verify slide 38 background (0.4 pts) ----------
    try:
        target_slide = prs.slides[TARGET_SLIDE_INDEX]
    except IndexError:
        print("✗ Slide 38 not present (0 pts)")
        return total_score  # early return – can’t earn further points

    fill = target_slide.background.fill
    if fill.type == MSO_FILL.GRADIENT:
        # Gather RGB colours of gradient stops
        stop_rgbs = {
            stop.color.rgb for stop in fill.gradient_stops if stop.color.rgb is not None
        }
        print(f"Slide 38 gradient stop colours: {sorted(stop_rgbs)}")
        if GRADIENT_BLUE2_COLORS.issubset(stop_rgbs):
            print("✓ Slide 38 uses Gradient Blue 2 (0.4 pts)")
            total_score += 0.4
        else:
            print("✗ Slide 38 gradient colours don’t match Blue 2 (0 pts)")
    else:
        print("✗ Slide 38 background is not a gradient (0 pts)")

    # ---------- 4. Ensure other slides remain plain (0.2 pts) ----------
    others_ok = True
    for idx, slide in enumerate(prs.slides):
        if idx == TARGET_SLIDE_INDEX:
            continue  # skip target slide – already checked
        if slide.background.fill.type == MSO_FILL.GRADIENT:
            print(f"✗ Slide {idx + 1} unexpectedly has a gradient background")
            others_ok = False
            break
    if others_ok:
        print("✓ All other slides keep original (non-gradient) background (0.2 pts)")
        total_score += 0.2

    # ---------- 5. Final score ----------
    final_score = min(total_score, MAX_SCORE)
    print(f"Total Score: {final_score} / {MAX_SCORE}")
    print(f"REWARD: {final_score}")
    return final_score


if __name__ == "__main__":
    verify_task(FILE_PATH)

