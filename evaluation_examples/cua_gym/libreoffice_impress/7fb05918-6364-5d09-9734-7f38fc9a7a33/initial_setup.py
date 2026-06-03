"""
Initial Setup: Create a 20-slide company report presentation with no custom slide shows.
Task ID: impress_fix_035
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_035'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
    return slide


def add_blank_with_text(prs, title_text, body_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title textbox
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
    # Body textbox
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = body_text
    return slide


def create_initial():
    prs = Presentation()

    # Slide 1: Title slide
    add_title_slide(prs, "Meridian Technologies — Annual Report 2025",
                    "Prepared by the Strategy & Finance Division")

    # Slide 2: Table of Contents
    add_content_slide(prs, "Table of Contents", [
        "1. Executive Overview",
        "2. Financial Performance",
        "3. Product Portfolio",
        "4. Market Analysis",
        "5. Regional Breakdown",
        "6. R&D Initiatives",
        "7. Customer Insights",
        "8. Operations & Supply Chain",
        "9. Human Resources",
        "10. Risk Assessment & Outlook",
    ])

    # Slide 3: Executive Overview
    add_content_slide(prs, "Executive Overview", [
        "Revenue grew 18.3% year-over-year to $2.47 billion",
        "Operating margin expanded to 23.1%, up from 19.8%",
        "Customer base surpassed 4.2 million active users globally",
        "Launched 3 new flagship products across enterprise and consumer segments",
        "Expanded operations to 12 new international markets",
    ])

    # Slide 4: Revenue Breakdown by Quarter
    add_content_slide(prs, "Revenue Breakdown by Quarter", [
        "Q1 2025: $542M (+12% YoY)",
        "Q2 2025: $598M (+16% YoY)",
        "Q3 2025: $631M (+19% YoY)",
        "Q4 2025: $699M (+24% YoY)",
        "Full year total: $2,470M",
    ])

    # Slide 5: Gross Margin Analysis
    add_content_slide(prs, "Gross Margin Analysis", [
        "Overall gross margin: 64.2% (up from 61.5%)",
        "Hardware segment: 42.8% margin",
        "Software segment: 81.3% margin",
        "Services segment: 55.7% margin",
        "Cost reduction initiatives saved $83M in COGS",
    ])

    # Slide 6: Operating Expenses
    add_content_slide(prs, "Operating Expenses", [
        "Total OpEx: $1,899M (76.9% of revenue)",
        "R&D: $618M (25.0% of revenue) — up from 22.1%",
        "Sales & Marketing: $494M (20.0%)",
        "General & Administrative: $247M (10.0%)",
        "Depreciation & Amortization: $148M",
    ])

    # Slide 7: Product Portfolio Highlights
    add_content_slide(prs, "Product Portfolio Highlights", [
        "Meridian CloudOS 5.0 — enterprise cloud platform (42% of revenue)",
        "Meridian Edge — IoT device management suite (18% of revenue)",
        "Meridian Analytics Pro — business intelligence tool (15% of revenue)",
        "Meridian Secure — zero-trust cybersecurity framework (12% of revenue)",
        "Meridian Connect — unified communications platform (8% of revenue)",
        "Other products and services (5% of revenue)",
    ])

    # Slide 8: Competitive Landscape
    add_content_slide(prs, "Competitive Landscape", [
        "Market share in enterprise cloud: 14.6% (#3 position)",
        "IoT management market: 22.1% (#1 position)",
        "Business intelligence market: 8.9% (#5 position)",
        "Key competitors: NovaTech, Apex Systems, Cirrus Digital",
        "Differentiation: integrated platform approach vs. point solutions",
    ])

    # Slide 9: Market Analysis — Addressable Markets
    add_content_slide(prs, "Market Analysis — Addressable Markets", [
        "Total addressable market (TAM): $186 billion by 2027",
        "Enterprise cloud services: $98B (growing at 19% CAGR)",
        "IoT management: $42B (growing at 24% CAGR)",
        "Cybersecurity: $28B (growing at 14% CAGR)",
        "Business analytics: $18B (growing at 11% CAGR)",
    ])

    # Slide 10: Geographic Revenue Distribution
    add_content_slide(prs, "Geographic Revenue Distribution", [
        "North America: $1,284M (52% of total)",
        "Europe, Middle East & Africa: $618M (25%)",
        "Asia-Pacific: $395M (16%)",
        "Latin America: $173M (7%)",
        "Fastest growing region: Asia-Pacific (+31% YoY)",
    ])

    # Slide 11: Key Customer Metrics
    add_content_slide(prs, "Key Customer Metrics", [
        "Total active customers: 4.2 million (+22% YoY)",
        "Enterprise customers (>$100K ARR): 2,847 accounts",
        "Net revenue retention rate: 118%",
        "Customer acquisition cost (CAC): $1,240 (down 8%)",
        "Average contract value (ACV): $28,500 (up 15%)",
    ])

    # Slide 12: Strategic Partnerships & Alliances
    add_content_slide(prs, "Strategic Partnerships & Alliances", [
        "Technology partnership with Samsung for IoT hardware integration",
        "Joint venture with Deutsche Telekom for European expansion",
        "Academic partnership with MIT for AI research collaboration",
        "Channel partnerships with 340+ managed service providers",
        "Government contracts worth $187M across 8 federal agencies",
    ])

    # Slide 13: R&D Investments & Innovation
    add_content_slide(prs, "R&D Investments & Innovation", [
        "R&D spend: $618M (25% of revenue)",
        "Headcount: 3,200 engineers across 7 global R&D centers",
        "Patent portfolio: 1,847 granted, 423 pending",
        "Key focus areas: AI/ML, quantum computing, edge intelligence",
        "Open source contributions: 28 projects, 12K+ GitHub stars",
    ])

    # Slide 14: Talent & Workforce
    add_content_slide(prs, "Talent & Workforce", [
        "Total headcount: 14,200 employees (+18% YoY)",
        "Engineering: 5,600 (39%)",
        "Sales & Marketing: 3,400 (24%)",
        "Operations: 2,800 (20%)",
        "G&A: 2,400 (17%)",
        "Employee satisfaction score: 4.3/5.0 (industry avg: 3.8)",
    ])

    # Slide 15: Forward Guidance & Outlook
    add_content_slide(prs, "Forward Guidance & Outlook", [
        "FY2026 revenue guidance: $2.9B–$3.1B (+17–25% YoY)",
        "Target operating margin: 25–27%",
        "Planned CapEx: $320M (data centers + R&D facilities)",
        "M&A pipeline: 2–3 tuck-in acquisitions totaling $200–400M",
        "Long-term goal: $5B revenue by 2028",
    ])

    # Slide 16: Risk Factors
    add_content_slide(prs, "Risk Factors", [
        "Macroeconomic uncertainty and potential recession impact",
        "Intensifying competition from hyperscalers (AWS, Azure, GCP)",
        "Regulatory changes in data privacy (EU AI Act, GDPR updates)",
        "Supply chain disruptions in semiconductor sourcing",
        "Talent retention in competitive labor market",
    ])

    # Slide 17: ESG & Sustainability Initiatives
    add_content_slide(prs, "ESG & Sustainability Initiatives", [
        "Carbon neutral operations achieved in Q3 2025",
        "100% renewable energy for all data centers by 2026",
        "Diversity metrics: 38% female leadership (up from 31%)",
        "Community investment: $12M in STEM education programs",
        "Supply chain audit: 94% compliance with ethical sourcing standards",
    ])

    # Slide 18: Capital Allocation & Shareholder Returns
    add_content_slide(prs, "Capital Allocation & Shareholder Returns", [
        "Free cash flow: $412M (16.7% FCF margin)",
        "Share buyback program: $200M authorized in Q2 2025",
        "Dividend increase: $0.48/share (up 12% from prior year)",
        "Cash and equivalents: $1.8B (strong balance sheet)",
        "Net debt/EBITDA ratio: 0.6x (well within investment grade)",
    ])

    # Slide 19: Key Performance Indicators Dashboard
    add_content_slide(prs, "Key Performance Indicators Dashboard", [
        "Revenue growth: 18.3% (target: 15%) ✓",
        "Operating margin: 23.1% (target: 22%) ✓",
        "Customer NPS: 72 (target: 65) ✓",
        "Employee engagement: 4.3/5.0 (target: 4.0) ✓",
        "Product uptime SLA: 99.97% (target: 99.95%) ✓",
    ])

    # Slide 20: Thank You / Contact
    add_title_slide(prs, "Thank You",
                    "Questions? Contact: investor.relations@meridiantech.com\n"
                    "Meridian Technologies Inc. | www.meridiantech.com")

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
