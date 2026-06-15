"""
Initial Setup: Import slides from Appendix.pptx into Report_Final.pptx
Task ID: impress_fix_040
Domain: libreoffice_impress

Creates:
  /home/user/Report_Final.pptx  - 12 slides with 'Corporate' master style
  /home/user/Desktop/Appendix.pptx - 8 slides with 'Simple_White' master style
Opens Report_Final.pptx in LibreOffice Impress.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_040'
OUTPUT_MAIN = f'{WORKDIR}/Report_Final.pptx'
OUTPUT_APPENDIX = f'{WORKDIR}/Desktop/Appendix.pptx'

# Corporate theme colors
CORP_DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
CORP_ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
CORP_LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
CORP_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CORP_TEXT_DARK = RGBColor(0x33, 0x33, 0x33)

# Simple White theme colors
SIMPLE_GRAY = RGBColor(0x66, 0x66, 0x66)
SIMPLE_BLACK = RGBColor(0x00, 0x00, 0x00)


def launch_gui(command: str, delay_sec: float = 1.0):
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_corporate_title_slide(prs, title_text, subtitle_text):
    """Add a title slide with Corporate styling."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = CORP_DARK_BLUE

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = CORP_WHITE

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(7), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle_text
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.name = "Calibri"
    run2.font.size = Pt(18)
    run2.font.color.rgb = CORP_ACCENT_BLUE
    return slide


def add_corporate_content_slide(prs, title_text, bullet_points):
    """Add a content slide with Corporate styling: dark blue header bar, bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Header bar
    header = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(1.2)  # MSO_SHAPE.RECTANGLE
    )
    header.fill.solid()
    header.fill.fore_color.rgb = CORP_DARK_BLUE
    header.line.fill.background()

    # Title text in header
    tf = header.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = CORP_WHITE

    # Accent line under header
    line = slide.shapes.add_shape(
        1, Inches(0), Inches(1.2), Inches(10), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CORP_ACCENT_BLUE
    line.line.fill.background()

    # Bullet content
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5))
    tf2 = txBox.text_frame
    tf2.word_wrap = True
    for i, bp in enumerate(bullet_points):
        if i == 0:
            p2 = tf2.paragraphs[0]
        else:
            p2 = tf2.add_paragraph()
        p2.text = bp
        p2.space_after = Pt(10)
        run2 = p2.runs[0]
        run2.font.name = "Calibri"
        run2.font.size = Pt(16)
        run2.font.color.rgb = CORP_TEXT_DARK
    return slide


def add_simple_white_slide(prs, title_text, bullet_points):
    """Add a content slide with Simple_White styling: plain, minimal."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title - simple black text, no background bar
    txBox_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.9))
    tf_t = txBox_title.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = title_text
    p_t.alignment = PP_ALIGN.LEFT
    run_t = p_t.runs[0]
    run_t.font.name = "Arial"
    run_t.font.size = Pt(22)
    run_t.font.bold = False
    run_t.font.color.rgb = SIMPLE_BLACK

    # Thin gray line under title
    line = slide.shapes.add_shape(
        1, Inches(0.8), Inches(1.4), Inches(8.4), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = SIMPLE_GRAY
    line.line.fill.background()

    # Content bullets - gray text, Arial
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(8.4), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bp in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bp
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(14)
        run.font.color.rgb = SIMPLE_GRAY
    return slide


def create_report_final():
    """Create the 12-slide Corporate-themed Report_Final.pptx."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_corporate_title_slide(prs,
        "Q4 2025 Strategic Review",
        "Northwind Dynamics Inc. | Board Presentation | December 2025")

    # Slide 2: Agenda
    add_corporate_content_slide(prs, "Agenda", [
        "1. Executive Summary",
        "2. Revenue Performance & Market Share",
        "3. Product Development Pipeline",
        "4. Customer Acquisition Metrics",
        "5. Regional Expansion Update",
        "6. Technology Infrastructure",
        "7. Risk Assessment & Mitigation",
        "8. Financial Outlook 2026",
        "9. Strategic Recommendations",
        "10. Q&A",
    ])

    # Slide 3: Executive Summary
    add_corporate_content_slide(prs, "Executive Summary", [
        "Total revenue reached $287.4M, exceeding target by 12.3%",
        "Customer base grew to 14,200 active accounts (+18% YoY)",
        "Launched 3 new product lines across enterprise and SMB segments",
        "EBITDA margin improved to 24.6%, up from 21.8% in Q3",
        "Successfully expanded into Southeast Asian markets",
    ])

    # Slide 4: Revenue Performance
    add_corporate_content_slide(prs, "Revenue Performance & Market Share", [
        "North America: $168.2M (+15.4% YoY) - 58.5% of total",
        "Europe: $72.1M (+9.8% YoY) - 25.1% of total",
        "Asia-Pacific: $35.8M (+32.1% YoY) - 12.5% of total",
        "Rest of World: $11.3M (+7.2% YoY) - 3.9% of total",
        "Market share increased from 18.2% to 21.7% in core segments",
        "Average deal size: $42,300 (up from $37,800)",
    ])

    # Slide 5: Product Development
    add_corporate_content_slide(prs, "Product Development Pipeline", [
        "CloudSync Enterprise v3.0 - Released Oct 15, adoption rate 67%",
        "DataVault Analytics Suite - Beta phase, 230 enterprise testers",
        "MobileFirst Platform 2.0 - Launch scheduled for Feb 2026",
        "AI-Powered Insights Module - Development phase, Q2 2026 target",
        "R&D investment: $41.2M (14.3% of revenue)",
    ])

    # Slide 6: Customer Metrics
    add_corporate_content_slide(prs, "Customer Acquisition Metrics", [
        "Net Promoter Score: 72 (industry benchmark: 54)",
        "Customer Acquisition Cost: $2,340 (down 8.1% from Q3)",
        "Customer Lifetime Value: $156,700 (up 11.2%)",
        "Churn rate: 2.8% (down from 3.4% in Q3)",
        "Enterprise segment retention: 97.2%",
        "Support ticket resolution: avg 4.2 hours (SLA target: 8 hours)",
    ])

    # Slide 7: Regional Expansion
    add_corporate_content_slide(prs, "Regional Expansion Update", [
        "Singapore office fully operational since September 2025",
        "Tokyo partnership with Yamada Corp signed, revenue sharing model",
        "Sydney satellite office hiring complete: 24 staff",
        "Berlin European HQ expansion approved for 2026",
        "Latin America feasibility study in progress (Brazil, Mexico focus)",
    ])

    # Slide 8: Technology
    add_corporate_content_slide(prs, "Technology Infrastructure", [
        "Cloud migration 94% complete (target: 100% by Q1 2026)",
        "Platform uptime: 99.97% (SLA: 99.95%)",
        "Security audit: zero critical vulnerabilities found",
        "API response time improved by 34% after CDN optimization",
        "Database performance: 12ms average query time (down from 28ms)",
    ])

    # Slide 9: Risk Assessment
    add_corporate_content_slide(prs, "Risk Assessment & Mitigation", [
        "Regulatory compliance: GDPR, SOC2, ISO 27001 all maintained",
        "Supply chain risk: reduced vendor concentration from 3 to 7 providers",
        "Talent retention: key person dependency mitigated with cross-training",
        "Currency exposure hedged for EUR, GBP, JPY through Q2 2026",
        "Cybersecurity insurance increased to $50M coverage",
    ])

    # Slide 10: Financial Outlook
    add_corporate_content_slide(prs, "Financial Outlook 2026", [
        "Projected revenue: $345-360M (20-25% growth)",
        "Target EBITDA margin: 26-28%",
        "Planned headcount increase: 180 positions (engineering, sales)",
        "Capital expenditure budget: $28.5M",
        "Dividend proposal: $1.85 per share (up from $1.60)",
    ])

    # Slide 11: Recommendations
    add_corporate_content_slide(prs, "Strategic Recommendations", [
        "Accelerate AI module development - allocate additional $8M R&D budget",
        "Pursue acquisition of DataStream Analytics (valuation: $45-55M)",
        "Expand Singapore team to serve as APAC regional hub",
        "Invest in partner channel program targeting 30% indirect revenue by 2027",
        "Launch customer advisory board with top 25 enterprise accounts",
    ])

    # Slide 12: Q&A / Closing
    add_corporate_title_slide(prs,
        "Thank You",
        "Questions & Discussion | Contact: strategy@northwind.com")

    prs.save(OUTPUT_MAIN)
    print(f"Created: {OUTPUT_MAIN} with {len(prs.slides)} slides")


def create_appendix():
    """Create the 8-slide Simple_White-themed Appendix.pptx."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Appendix Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Appendix"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(32)
    run.font.color.rgb = SIMPLE_BLACK

    p2 = tf.add_paragraph()
    p2.text = "Supporting Data & Detailed Analysis"
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.name = "Arial"
    run2.font.size = Pt(16)
    run2.font.color.rgb = SIMPLE_GRAY

    # Slide 2: Detailed Revenue Breakdown
    add_simple_white_slide(prs, "A1: Detailed Revenue Breakdown by Product Line", [
        "CloudSync Enterprise: $124.8M (43.4% of total)",
        "DataVault Standard: $68.3M (23.8%)",
        "MobileFirst Platform: $52.1M (18.1%)",
        "Professional Services: $28.6M (10.0%)",
        "Support & Maintenance: $13.6M (4.7%)",
        "YoY growth highest in MobileFirst (+41.2%)",
    ])

    # Slide 3: Employee Headcount
    add_simple_white_slide(prs, "A2: Employee Headcount & Distribution", [
        "Total headcount: 1,247 (up from 1,089 in Q4 2024)",
        "Engineering: 498 (39.9%)",
        "Sales & Marketing: 312 (25.0%)",
        "Customer Success: 187 (15.0%)",
        "G&A: 156 (12.5%)",
        "Executive & Strategy: 94 (7.5%)",
        "Average tenure: 3.4 years | Voluntary turnover: 8.2%",
    ])

    # Slide 4: Customer Segmentation
    add_simple_white_slide(prs, "A3: Customer Segmentation Analysis", [
        "Enterprise (>$100K ARR): 412 accounts, $198.4M revenue",
        "Mid-Market ($25K-$100K): 1,834 accounts, $62.7M revenue",
        "SMB (<$25K): 11,954 accounts, $26.3M revenue",
        "Top 10 accounts represent 14.2% of total revenue",
        "Concentration risk within acceptable thresholds",
    ])

    # Slide 5: Competitive Landscape
    add_simple_white_slide(prs, "A4: Competitive Landscape", [
        "Primary competitor TechFlow Inc: 28.4% market share (down from 30.1%)",
        "CloudBase Corp: 15.8% market share (stable)",
        "DataPrime Solutions: 12.3% (new entrant, growing rapidly)",
        "Northwind win rate vs. TechFlow: 62% in head-to-head deals",
        "Key differentiator: integrated analytics + superior support SLAs",
    ])

    # Slide 6: Technology Stack Details
    add_simple_white_slide(prs, "A5: Technology Stack & Architecture", [
        "Primary cloud: AWS (68%), Azure (22%), GCP (10%)",
        "Microservices count: 247 (up from 198)",
        "CI/CD pipeline: avg 12.4 deployments/day",
        "Test coverage: 87.3% (target: 90% by Q2 2026)",
        "Incident response time: P1 avg 8 minutes, P2 avg 23 minutes",
        "Tech debt reduction: 34% of sprint capacity allocated",
    ])

    # Slide 7: Partnership Pipeline
    add_simple_white_slide(prs, "A6: Partnership & Integration Pipeline", [
        "Active technology partners: 42 (up from 31)",
        "Salesforce integration: GA, 1,200+ installations",
        "ServiceNow connector: beta, 89 pilot customers",
        "SAP integration: development phase, Q1 2026 target",
        "Microsoft Teams app: 4,500 monthly active users",
        "Partner-sourced revenue: $34.2M (11.9% of total)",
    ])

    # Slide 8: Glossary & Definitions
    add_simple_white_slide(prs, "A7: Glossary & Key Definitions", [
        "ARR: Annual Recurring Revenue",
        "EBITDA: Earnings Before Interest, Taxes, Depreciation, Amortization",
        "NPS: Net Promoter Score (range -100 to +100)",
        "CAC: Customer Acquisition Cost",
        "LTV: Customer Lifetime Value",
        "SLA: Service Level Agreement",
        "YoY: Year-over-Year comparison",
    ])

    os.makedirs(f'{WORKDIR}/Desktop', exist_ok=True)
    prs.save(OUTPUT_APPENDIX)
    print(f"Created: {OUTPUT_APPENDIX} with {len(prs.slides)} slides")


def main():
    create_report_final()
    create_appendix()

    # Open Report_Final.pptx in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT_MAIN}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


main()
