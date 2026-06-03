"""
Reward Script: Import slides from Appendix.pptx into Report_Final.pptx
Task ID: impress_fix_040
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Slide count == 20
  Component 2 (0.30): Appended slides 13-20 contain correct content from Appendix
  Component 3 (0.20): Original slides 1-12 content preserved
  Component 4 (0.20): All slides use a single slide master (unified style)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_040'

# Expected title texts for the 8 appended slides (from Appendix.pptx)
APPENDIX_TITLE_KEYWORDS = [
    "Appendix",                                    # Slide 13
    "A1: Detailed Revenue Breakdown",              # Slide 14
    "A2: Employee Headcount",                      # Slide 15
    "A3: Customer Segmentation",                   # Slide 16
    "A4: Competitive Landscape",                   # Slide 17
    "A5: Technology Stack",                         # Slide 18
    "A6: Partnership & Integration",               # Slide 19
    "A7: Glossary & Key Definitions",              # Slide 20
]

# Expected title keywords for the original 12 slides
ORIGINAL_TITLE_KEYWORDS = [
    "Q4 2025 Strategic Review",
    "Agenda",
    "Executive Summary",
    "Revenue Performance",
    "Product Development",
    "Customer Acquisition",
    "Regional Expansion",
    "Technology Infrastructure",
    "Risk Assessment",
    "Financial Outlook",
    "Strategic Recommendations",
    "Thank You",
]


def get_slide_texts(slide):
    """Get all text from a slide's shapes."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                texts.append(t)
    return texts


def slide_contains_keyword(slide, keyword):
    """Check if any text in the slide contains the keyword."""
    for text in get_slide_texts(slide):
        if keyword.lower() in text.lower():
            return True
    return False


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Slide count is exactly 20 (0.30 points)
    try:
        if num_slides == 20:
            print(f"PASS: Component 1 — Slide count is 20 (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 1 — Expected 20 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Appended slides 13-20 contain Appendix content (0.30 points)
    # Each of the 8 slides contributes 0.30/8 = 0.0375 points
    try:
        if num_slides >= 20:
            matched = 0
            for i, keyword in enumerate(APPENDIX_TITLE_KEYWORDS):
                slide_idx = 12 + i  # slides 13-20 are indices 12-19
                slide = prs.slides[slide_idx]
                if slide_contains_keyword(slide, keyword):
                    matched += 1
                else:
                    actual_texts = get_slide_texts(slide)
                    print(f"  FAIL: Slide {slide_idx+1} missing keyword '{keyword}', found: {[t[:50] for t in actual_texts]}")

            points = 0.30 * (matched / 8)
            if matched == 8:
                print(f"PASS: Component 2 — All 8 appended slides have correct content ({points:.2f} pts)")
            elif matched > 0:
                print(f"PARTIAL: Component 2 — {matched}/8 appended slides matched ({points:.2f} pts)")
            else:
                print(f"FAIL: Component 2 — 0/8 appended slides matched")
            if points > 0:
                total_score += points
        else:
            print(f"FAIL: Component 2 — Not enough slides ({num_slides}) to check appended content")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Original slides 1-12 content preserved AND slides were appended (0.20 points)
    # This is a compound check: must have 20 slides (task change) AND originals intact
    # The 20-slide gate ensures this only passes AFTER the task is done
    try:
        if num_slides == 20:
            matched = 0
            for i, keyword in enumerate(ORIGINAL_TITLE_KEYWORDS):
                slide = prs.slides[i]
                if slide_contains_keyword(slide, keyword):
                    matched += 1
                else:
                    actual_texts = get_slide_texts(slide)
                    print(f"  FAIL: Original slide {i+1} missing keyword '{keyword}', found: {[t[:50] for t in actual_texts]}")

            points = 0.20 * (matched / 12)
            if matched == 12:
                print(f"PASS: Component 3 — All 12 original slides preserved with 20 total ({points:.2f} pts)")
            elif matched > 0:
                print(f"PARTIAL: Component 3 — {matched}/12 original slides matched ({points:.2f} pts)")
            else:
                print(f"FAIL: Component 3 — 0/12 original slides matched")
            if points > 0:
                total_score += points
        else:
            print(f"FAIL: Component 3 — Slide count is {num_slides}, not 20; cannot confirm originals preserved after merge")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: All slides use a single slide master (0.20 points)
    # This verifies the "adopt current presentation's master slide style" requirement
    try:
        if num_slides >= 20:
            master_ids = set()
            for slide in prs.slides:
                master_ids.add(id(slide.slide_layout.slide_master))

            num_masters_in_file = len(prs.slide_masters)

            if len(master_ids) == 1 and num_masters_in_file == 1:
                print(f"PASS: Component 4 — All slides use single master, 1 master in file (0.20 pts)")
                total_score += 0.20
            elif len(master_ids) == 1:
                print(f"PARTIAL: Component 4 — All slides reference same master but {num_masters_in_file} masters in file (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 4 — {len(master_ids)} different masters used across slides, {num_masters_in_file} masters in file")
        else:
            print(f"FAIL: Component 4 — Not enough slides to check master unification")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Main entry point
file_path = f'{WORKDIR}/Report_Final.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
