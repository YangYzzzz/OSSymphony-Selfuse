"""
Reward Script: Verify 'Two Content' master slide layout creation
Task ID: impress_gf3_028
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): 'Two Content' layout exists by name
  Component 2 (0.25): Title placeholder in top ~25% of slide
  Component 3 (0.25): Two OBJECT placeholders (content areas) exist below title
  Component 4 (0.25): Content placeholders are equal-width and side by side
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_028'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Find the 'Two Content' layout across all masters
    two_content_layout = None
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name.strip().lower() == 'two content':
                two_content_layout = layout
                break
        if two_content_layout is not None:
            break

    # Component 1: 'Two Content' layout exists (0.25 points)
    try:
        if two_content_layout is not None:
            print(f"PASS: Component 1 — 'Two Content' layout found (name='{two_content_layout.name}') (0.25 pts)")
            total_score += 0.25
        else:
            layout_names = []
            for master in prs.slide_masters:
                for layout in master.slide_layouts:
                    layout_names.append(layout.name)
            print(f"FAIL: Component 1 — No 'Two Content' layout found. Existing layouts: {layout_names}")
            # If the layout doesn't exist, nothing else can pass
            print(f"\nScore: {total_score}/1.0")
            print(f"REWARD: {total_score}")
            return total_score
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Gather placeholders from the layout, excluding footer/date/slide-number types
    # (placeholder idx >= 10 are typically footer/date/slidenum)
    content_placeholders = []
    title_placeholder = None
    for ph in two_content_layout.placeholders:
        ph_type = ph.placeholder_format.type
        ph_idx = ph.placeholder_format.idx
        # Skip date, footer, slide number placeholders
        if ph_idx >= 10:
            continue
        if ph_type is not None and str(ph_type) == 'TITLE (1)':
            title_placeholder = ph
        elif ph_type is not None and str(ph_type) == 'OBJECT (7)':
            content_placeholders.append(ph)

    # Component 2: Title placeholder in top ~25% of slide (0.25 points)
    try:
        if title_placeholder is not None:
            title_bottom = title_placeholder.top + title_placeholder.height
            # Title should be in roughly the top 25% of the slide
            # Allow some tolerance: title bottom should be within top 35% of slide height
            title_bottom_pct = title_bottom / slide_height
            title_width_pct = title_placeholder.width / slide_width

            if title_bottom_pct <= 0.35 and title_width_pct >= 0.70:
                print(f"PASS: Component 2 — Title at top (bottom={title_bottom_pct:.1%}, width={title_width_pct:.1%}) (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 2 — Title placement out of range "
                      f"(bottom at {title_bottom_pct:.1%}, width={title_width_pct:.1%})")
        else:
            print("FAIL: Component 2 — No TITLE placeholder found in 'Two Content' layout")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Two OBJECT (content) placeholders exist below title (0.25 points)
    try:
        if len(content_placeholders) >= 2:
            # Both content placeholders should be below the title
            if title_placeholder is not None:
                title_bottom = title_placeholder.top + title_placeholder.height
            else:
                title_bottom = 0

            below_title = [ph for ph in content_placeholders if ph.top >= title_bottom * 0.85]
            if len(below_title) >= 2:
                print(f"PASS: Component 3 — {len(below_title)} content placeholders below title (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 3 — Only {len(below_title)} content placeholders below title "
                      f"(need 2). Tops: {[ph.top for ph in content_placeholders]}, title_bottom: {title_bottom}")
        else:
            print(f"FAIL: Component 3 — Only {len(content_placeholders)} OBJECT placeholders found (need 2)")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Content placeholders are equal-width and side by side (0.25 points)
    try:
        if len(content_placeholders) >= 2:
            # Sort by left position to get left and right
            sorted_phs = sorted(content_placeholders, key=lambda ph: ph.left)
            left_ph = sorted_phs[0]
            right_ph = sorted_phs[1]

            # Check equal width (within 5% tolerance)
            width_ratio = min(left_ph.width, right_ph.width) / max(left_ph.width, right_ph.width)
            # Check they are side by side: right starts after left ends (or with small gap)
            gap = right_ph.left - (left_ph.left + left_ph.width)
            # Gap should be positive (side by side) and not too large (< 15% of slide width)
            gap_pct = gap / slide_width
            # Check tops are approximately equal (within 5% of slide height)
            top_diff_pct = abs(left_ph.top - right_ph.top) / slide_height

            all_ok = (width_ratio >= 0.85 and gap >= 0 and gap_pct <= 0.15 and top_diff_pct <= 0.05)
            details = []

            details.append(f"width ratio={width_ratio:.2f}")
            details.append(f"gap={gap_pct:.1%} of slide width")
            details.append(f"tops diff={top_diff_pct:.1%}")

            if all_ok:
                print(f"PASS: Component 4 — Content areas side by side, equal width ({'; '.join(details)}) (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 4 — {'; '.join(details)}")
        else:
            print("FAIL: Component 4 — Not enough content placeholders to verify layout")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
