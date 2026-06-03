"""
Initial Setup: Create a presentation with 8 slides, only 'Title and Content' and 'Blank' layouts.
Task ID: impress_gf3_028
Domain: libreoffice_impress
"""

import os
import re
import shlex
import subprocess
import time
import copy
import zipfile
import io
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_028'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def remove_unused_layouts(pptx_path, keep_names):
    """
    Remove slide layouts not in keep_names from the pptx file.
    This operates at the ZIP/XML level to properly remove layout files and references.
    """
    from pptx import Presentation as Prs
    # First, figure out which layout files to keep
    prs = Prs(pptx_path)
    keep_indices = []
    for i, layout in enumerate(prs.slide_layouts):
        if layout.name in keep_names:
            keep_indices.append(i)

    # Work at ZIP level
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zout:
            # Determine which slideLayout files to keep
            # Parse the slideMaster1.xml.rels to find layout rIds
            master_rels_path = 'ppt/slideMasters/_rels/slideMaster1.xml.rels'
            master_rels_xml = zin.read(master_rels_path)
            rels_root = etree.fromstring(master_rels_xml)

            # Find all slideLayout relationships
            layout_rels = []
            for rel in rels_root:
                target = rel.get('Target')
                if target and 'slideLayout' in target:
                    layout_rels.append(rel)

            # Sort by layout number
            layout_rels.sort(key=lambda r: int(re.search(r'slideLayout(\d+)\.xml', r.get('Target')).group(1)))

            # Determine which layout files to remove
            remove_layout_files = set()  # e.g. 'ppt/slideLayouts/slideLayout3.xml'
            remove_layout_basenames = set()  # e.g. 'slideLayout3.xml'
            remove_rids = set()
            for i, rel in enumerate(layout_rels):
                if i not in keep_indices:
                    target = rel.get('Target')
                    # target is like ../slideLayouts/slideLayoutN.xml
                    basename = target.split('/')[-1]  # slideLayoutN.xml
                    full_path = f'ppt/slideLayouts/{basename}'
                    remove_layout_files.add(full_path)
                    remove_layout_basenames.add(basename)
                    remove_rids.add(rel.get('Id'))

            # Copy all files, skipping removed layouts and their rels
            for item in zin.namelist():
                norm = item.replace('\\', '/')
                skip = False
                if norm in remove_layout_files:
                    skip = True
                for bn in remove_layout_basenames:
                    if norm == f'ppt/slideLayouts/_rels/{bn}.rels':
                        skip = True
                        break

                if skip:
                    continue

                data = zin.read(item)

                # Fix slideMaster1.xml.rels - remove refs to deleted layouts
                if norm == master_rels_path:
                    root = etree.fromstring(data)
                    for rel in list(root):
                        if rel.get('Id') in remove_rids:
                            root.remove(rel)
                    data = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)

                # Fix slideMaster1.xml - remove sldLayoutId refs
                if norm == 'ppt/slideMasters/slideMaster1.xml':
                    root = etree.fromstring(data)
                    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                          'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}
                    sld_layout_id_lst = root.find('.//p:sldLayoutIdLst', ns)
                    if sld_layout_id_lst is not None:
                        for child in list(sld_layout_id_lst):
                            rid = child.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                            if rid in remove_rids:
                                sld_layout_id_lst.remove(child)
                    data = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)

                # Fix [Content_Types].xml - remove references to deleted layout files
                if norm == '[Content_Types].xml':
                    root = etree.fromstring(data)
                    for child in list(root):
                        part_name = child.get('PartName', '')
                        for rf in remove_layout_files:
                            if part_name == '/' + rf:
                                root.remove(child)
                                break
                    data = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)

                zout.writestr(item, data)

    with open(pptx_path, 'wb') as f:
        f.write(buf.getvalue())


def create_initial():
    prs = Presentation()

    # Layout indices in default template:
    # 1 = Title and Content
    # 6 = Blank
    layout_content = prs.slide_layouts[1]  # Title and Content
    layout_blank = prs.slide_layouts[6]    # Blank

    # --- Slide 1: Title slide (using Title and Content layout) ---
    slide1 = prs.slides.add_slide(layout_content)
    slide1.shapes.title.text = "Q4 2025 Strategic Planning"
    slide1.placeholders[1].text = (
        "Prepared by the Operations & Strategy Team\n"
        "Meridian Consulting Group\n"
        "November 2025"
    )

    # --- Slide 2: Market Overview ---
    slide2 = prs.slides.add_slide(layout_content)
    slide2.shapes.title.text = "Market Overview"
    tf2 = slide2.placeholders[1].text_frame
    tf2.clear()
    bullets = [
        "Total addressable market grew 12.3% YoY to $47.8B",
        "North America remains dominant at 38% market share",
        "APAC region showing fastest growth trajectory at 23% CAGR",
        "Key competitors: Horizon Labs, NovaTech Solutions, Atlas Corp",
        "Regulatory environment tightening in EU markets",
    ]
    for i, txt in enumerate(bullets):
        if i == 0:
            tf2.paragraphs[0].text = txt
        else:
            p = tf2.add_paragraph()
            p.text = txt

    # --- Slide 3: Revenue Breakdown ---
    slide3 = prs.slides.add_slide(layout_content)
    slide3.shapes.title.text = "Revenue Breakdown by Division"
    tf3 = slide3.placeholders[1].text_frame
    tf3.clear()
    divisions = [
        "Enterprise Solutions: $18.4M (42% of total)",
        "Cloud Infrastructure: $12.7M (29% of total)",
        "Professional Services: $8.1M (18% of total)",
        "Product Licensing: $4.8M (11% of total)",
        "Total Revenue: $44.0M | Target: $48.5M",
    ]
    for i, txt in enumerate(divisions):
        if i == 0:
            tf3.paragraphs[0].text = txt
        else:
            p = tf3.add_paragraph()
            p.text = txt

    # --- Slide 4: Team Structure (blank layout with textbox) ---
    slide4 = prs.slides.add_slide(layout_blank)
    txBox = slide4.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Organizational Structure"
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    txBox2 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    lines = [
        "CEO: Rebecca Thornton",
        "  VP Engineering: David Park (32 reports)",
        "  VP Sales: Amanda Foster (28 reports)",
        "  VP Marketing: Carlos Rivera (15 reports)",
        "  VP Finance: Sarah Mitchell (12 reports)",
        "  VP Operations: James Okonkwo (22 reports)",
        "",
        "Total headcount: 247 (up from 198 in Q3)",
        "Open positions: 18 across Engineering and Sales",
    ]
    for i, line in enumerate(lines):
        if i == 0:
            tf2.paragraphs[0].text = line
        else:
            p = tf2.add_paragraph()
            p.text = line

    # --- Slide 5: Key Metrics ---
    slide5 = prs.slides.add_slide(layout_content)
    slide5.shapes.title.text = "Key Performance Metrics"
    tf5 = slide5.placeholders[1].text_frame
    tf5.clear()
    metrics = [
        "Customer Acquisition Cost: $2,340 (down 8%)",
        "Monthly Recurring Revenue: $3.67M",
        "Net Promoter Score: 72 (industry avg: 58)",
        "Customer Retention Rate: 94.2%",
        "Average Deal Size: $145K (up 12% QoQ)",
        "Pipeline Coverage: 3.8x (target: 3.0x)",
    ]
    for i, txt in enumerate(metrics):
        if i == 0:
            tf5.paragraphs[0].text = txt
        else:
            p = tf5.add_paragraph()
            p.text = txt

    # --- Slide 6: Regional Performance (blank with table-like textboxes) ---
    slide6 = prs.slides.add_slide(layout_blank)
    txTitle = slide6.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(0.8))
    tf = txTitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Regional Performance Summary"
    p.font.size = Pt(28)
    p.font.bold = True

    # Add table
    rows, cols = 6, 4
    tbl_shape = slide6.shapes.add_table(rows, cols, Inches(0.8), Inches(1.5), Inches(8), Inches(4))
    table = tbl_shape.table
    headers = ["Region", "Revenue", "Growth", "Target"]
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
    data_rows = [
        ["North America", "$18.2M", "+14%", "$17.5M"],
        ["Europe", "$11.4M", "+8%", "$12.0M"],
        ["Asia Pacific", "$8.7M", "+23%", "$8.0M"],
        ["Latin America", "$3.4M", "+11%", "$3.5M"],
        ["Middle East & Africa", "$2.3M", "+6%", "$2.5M"],
    ]
    for r, row_data in enumerate(data_rows, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 7: Strategic Initiatives ---
    slide7 = prs.slides.add_slide(layout_content)
    slide7.shapes.title.text = "Strategic Initiatives for Q1 2026"
    tf7 = slide7.placeholders[1].text_frame
    tf7.clear()
    initiatives = [
        "Launch AI-powered analytics dashboard (Project Orion)",
        "Expand APAC sales team by 40% with Singapore hub",
        "Migrate legacy customers to cloud platform (Phase 2)",
        "Achieve SOC 2 Type II certification by March",
        "Partner integration with Salesforce and HubSpot",
        "Reduce infrastructure costs by 15% via containerization",
    ]
    for i, txt in enumerate(initiatives):
        if i == 0:
            tf7.paragraphs[0].text = txt
        else:
            p = tf7.add_paragraph()
            p.text = txt

    # --- Slide 8: Next Steps ---
    slide8 = prs.slides.add_slide(layout_content)
    slide8.shapes.title.text = "Next Steps & Action Items"
    tf8 = slide8.placeholders[1].text_frame
    tf8.clear()
    steps = [
        "Finalize budget allocation by November 30",
        "Schedule department head reviews (Dec 1-15)",
        "Submit board presentation draft by December 20",
        "Complete hiring pipeline for Q1 positions",
        "Distribute updated OKRs to all team leads",
    ]
    for i, txt in enumerate(steps):
        if i == 0:
            tf8.paragraphs[0].text = txt
        else:
            p = tf8.add_paragraph()
            p.text = txt

    prs.save(OUTPUT)
    print(f'Saved initial presentation: {OUTPUT}')

    # Remove all layouts except "Title and Content" and "Blank"
    remove_unused_layouts(OUTPUT, keep_names={"Title and Content", "Blank"})
    print('Removed unused layouts, keeping only Title and Content + Blank')

    # Verify
    prs2 = Presentation(OUTPUT)
    print(f'Slide count: {len(prs2.slides)}')
    print('Available layouts:')
    for i, layout in enumerate(prs2.slide_layouts):
        print(f'  {i}: {layout.name}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
