"""
Initial Setup: Scientific Method presentation with bullet points on slide 3
Task ID: impress_stu_081
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_081'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "The Scientific Method"
    slide1.placeholders[1].text = "A Step-by-Step Guide to Scientific Inquiry"

    # --- Slide 2: Introduction ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "What Is the Scientific Method?"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "The scientific method is a systematic approach used by researchers worldwide to investigate phenomena, acquire new knowledge, and correct previous understanding."
    p2 = body2.add_paragraph()
    p2.text = "It provides a structured framework that minimizes bias and ensures reproducibility across experiments."
    p2a = body2.add_paragraph()
    p2a.text = ""
    p2b = body2.add_paragraph()
    p2b.text = "Developed over centuries, this method forms the backbone of modern scientific discovery."

    # --- Slide 3: Scientific Method Steps (6 bullet points) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Steps of the Scientific Method"
    body3 = slide3.placeholders[1].text_frame

    steps = [
        "Make an observation about a natural phenomenon",
        "Formulate a testable hypothesis",
        "Design and conduct an experiment",
        "Collect and analyze the data",
        "Draw conclusions from the results",
        "Communicate findings and repeat if necessary",
    ]

    # First paragraph uses existing paragraph
    body3.text = steps[0]
    for step_text in steps[1:]:
        p = body3.add_paragraph()
        p.text = step_text

    # Ensure all paragraphs have default bullet style (level 0)
    for para in body3.paragraphs:
        para.level = 0
        for run in para.runs:
            run.font.size = Pt(18)

    # --- Slide 4: Practical Applications ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Applications in Modern Research"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Clinical trials in pharmaceutical development follow the scientific method rigorously to test new treatments."
    p4a = body4.add_paragraph()
    p4a.text = "Environmental scientists use it to study climate change patterns and predict future trends."
    p4b = body4.add_paragraph()
    p4b.text = "Engineers apply the method when designing and testing new materials for aerospace applications."
    p4c = body4.add_paragraph()
    p4c.text = "Social scientists adapt the framework to study human behavior and societal trends."

    # --- Slide 5: Summary ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Key Takeaways"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "The scientific method provides an objective, repeatable process for investigating questions."
    p5a = body5.add_paragraph()
    p5a.text = "Following each step carefully ensures reliable and valid results."
    p5b = body5.add_paragraph()
    p5b.text = "Critical thinking and skepticism are essential at every stage of the process."
    p5c = body5.add_paragraph()
    p5c.text = "Science progresses through continuous cycles of hypothesis testing and refinement."

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
