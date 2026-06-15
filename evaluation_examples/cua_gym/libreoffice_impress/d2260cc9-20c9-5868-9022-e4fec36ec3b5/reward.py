"""
Reward script for impress_gf5_038:
Verify conditional row shading on the table in slide 5 of board_presentation.pptx.
  - 'At Risk'   -> orange (#FFA500)
  - 'On Track'  -> green  (#90EE90)
  - 'Completed' -> gray   (#D3D3D3)
8 data rows (rows 1-8, 0-indexed). Header row 0 is not checked.
"""

import subprocess, time, os

# --- Step 1.8: Persistence hook for LibreOffice ---
# If LibreOffice has the file open, Ctrl+S saves pending changes to disk.
try:
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        ["python3", "-c",
         "import pyautogui, time; time.sleep(0.3); pyautogui.hotkey('ctrl','s'); time.sleep(1.0)"],
        env=env, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
    )
    time.sleep(2.0)
except Exception:
    pass

# --- Verification ---
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

FILE_PATH = "/home/user/impress_gf5_038.pptx"
STATUS_COL = 3  # 0-indexed column for "Status"

COLOR_MAP = {
    "At Risk":   "FFA500",
    "On Track":  "90EE90",
    "Completed": "D3D3D3",
}

def get_cell_fill_rgb(cell):
    """Return hex RGB string of cell solid fill, or None."""
    try:
        fill = cell.fill
        if fill.type is not None:
            return str(fill.fore_color.rgb).upper()
    except Exception:
        pass
    return None

def main():
    try:
        prs = Presentation(FILE_PATH)
    except Exception as e:
        print(f"Cannot open file: {e}")
        print("REWARD: 0.0")
        return

    # Locate slide 5 (0-indexed = 4)
    if len(prs.slides) < 5:
        print("Slide 5 does not exist.")
        print("REWARD: 0.0")
        return

    slide = prs.slides[4]

    # Find the table
    table = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            break

    if table is None:
        print("No table found on slide 5.")
        print("REWARD: 0.0")
        return

    num_rows = len(table.rows)
    num_cols = len(table.columns)
    print(f"Table found: {num_rows} rows x {num_cols} cols")

    if num_rows < 9:
        print(f"Expected at least 9 rows (1 header + 8 data), got {num_rows}")
        print("REWARD: 0.0")
        return

    # Score each data row (rows 1-8)
    total_data_rows = 8
    correct_rows = 0

    for row_idx in range(1, total_data_rows + 1):
        status_text = table.cell(row_idx, STATUS_COL).text.strip()
        expected_color = COLOR_MAP.get(status_text)

        if expected_color is None:
            print(f"  Row {row_idx}: Unknown status '{status_text}' - skip (no credit)")
            continue

        row_ok = True
        for col_idx in range(num_cols):
            cell = table.cell(row_idx, col_idx)
            actual_color = get_cell_fill_rgb(cell)
            if actual_color != expected_color:
                print(f"  Row {row_idx}, Col {col_idx}: expected {expected_color}, got {actual_color}")
                row_ok = False
                break

        if row_ok:
            print(f"  Row {row_idx} ('{status_text}'): CORRECT ({expected_color})")
            correct_rows += 1
        else:
            print(f"  Row {row_idx} ('{status_text}'): INCORRECT")

    score = round(correct_rows / total_data_rows, 2)
    print(f"\nCorrect rows: {correct_rows}/{total_data_rows}")
    print(f"REWARD: {score}")

if __name__ == "__main__":
    main()
