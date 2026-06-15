"""initial_setup.py - Create board_presentation.pptx with 5 slides, table on slide 5 (no row shading)."""

import os
import subprocess
import shlex
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

FILE_PATH = "/home/user/impress_gf5_038.pptx"

prs = Presentation()

# ── Slide 1: Title Slide ──
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
slide1.shapes.title.text = "Board of Directors Quarterly Review"
slide1.placeholders[1].text = "Q1 2026 — Strategic Initiatives & Portfolio Health"

# ── Slide 2: Agenda ──
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
slide2.shapes.title.text = "Agenda"
tf = slide2.placeholders[1].text_frame
tf.text = "1. Financial Summary"
for item in ["2. Strategic Initiative Updates", "3. Risk Assessment", "4. Resource Allocation", "5. Project Portfolio Status"]:
    p = tf.add_paragraph()
    p.text = item

# ── Slide 3: Financial Summary ──
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
slide3.shapes.title.text = "Financial Summary"
tf3 = slide3.placeholders[1].text_frame
tf3.text = "Revenue: $42.3M (+12% YoY)"
for line in ["Operating Margin: 18.5%", "Cash Position: $85.2M", "Capital Expenditure: $6.1M"]:
    p = tf3.add_paragraph()
    p.text = line

# ── Slide 4: Risk Assessment ──
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
slide4.shapes.title.text = "Risk Assessment Overview"
tf4 = slide4.placeholders[1].text_frame
tf4.text = "Supply chain disruptions remain the top risk factor."
for line in [
    "Mitigation plans in place for 3 critical vendors.",
    "Cybersecurity audit completed — no critical findings.",
    "Regulatory compliance review scheduled for Q2.",
]:
    p = tf4.add_paragraph()
    p.text = line

# ── Slide 5: Project Portfolio Table ──
slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

# Add a title text box
title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.6))
tp = title_box.text_frame.paragraphs[0]
tp.text = "Project Portfolio Status"
tp.alignment = PP_ALIGN.CENTER
run = tp.runs[0]
run.font.size = Pt(24)
run.font.bold = True

# Table data: 1 header + 8 data rows, 5 columns
headers = ["Project", "Owner", "Budget", "Status", "Due Date"]
data_rows = [
    ["Cloud Migration",       "Sarah Chen",    "$2.4M",  "On Track",   "2026-06-30"],
    ["ERP Upgrade",           "James Wilson",  "$5.1M",  "At Risk",    "2026-09-15"],
    ["Mobile App Redesign",   "Priya Patel",   "$1.8M",  "On Track",   "2026-05-20"],
    ["Data Warehouse",        "Michael Torres", "$3.2M", "At Risk",    "2026-08-01"],
    ["Security Overhaul",     "Lisa Kim",      "$2.9M",  "Completed",  "2026-03-15"],
    ["Customer Portal",       "David Brown",   "$1.5M",  "On Track",   "2026-07-10"],
    ["AI Analytics Platform", "Rachel Green",  "$4.0M",  "At Risk",    "2026-10-30"],
    ["Office Relocation",     "Tom Harris",    "$0.8M",  "Completed",  "2026-02-28"],
]

rows = 9  # 1 header + 8 data
cols = 5
table_shape = slide5.shapes.add_table(rows, cols, Inches(0.3), Inches(1.0), Inches(9.0), Inches(5.5))
table = table_shape.table

# Set column widths
col_widths = [Inches(2.2), Inches(1.8), Inches(1.2), Inches(1.5), Inches(1.5)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w

# Header row styling
for c, header in enumerate(headers):
    cell = table.cell(0, c)
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = header
    run.font.bold = True
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # Dark blue header background
    cell_fill = cell.fill
    cell_fill.solid()
    cell_fill.fore_color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

# Data rows — NO shading applied
for r, row_data in enumerate(data_rows, start=1):
    for c, value in enumerate(row_data):
        cell = table.cell(r, c)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if c >= 2 else PP_ALIGN.LEFT
        run = p.add_run()
        run.text = value
        run.font.size = Pt(11)

prs.save(FILE_PATH)
print(f"Saved initial presentation to {FILE_PATH}")

# Launch LibreOffice Impress
def launch_gui(command, delay_sec=2.0):
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)

launch_gui(f'libreoffice --impress "{FILE_PATH}"', delay_sec=2.0)
print("LibreOffice Impress launched.")
