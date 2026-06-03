"""
FINAL REWARD SCRIPT - SUCCESS
Task: Anchor the second image To Paragraph and position it on the left.
Generated: 2025-10-17 12:30:28
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

"""
Reward Script: Anchor the second image To Paragraph and position it on the left
--------------------------------------------------------------------------
Verification Logic:
1. Make sure the presentation file exists and can be opened (no points – prerequisite).
2. Count the number of picture shapes across all slides.
   • Award 0.3 points ONLY if at least two images exist (task could be attempted).
3. Identify the *second* image in the document order (the order they are stored in
   the PPTX file).  Verify its horizontal (left) position:
   • Slide width is available via `prs.slide_width` (EMU units).
   • If the picture's left coordinate is within 5 % of the slide’s left edge we
     deem it “positioned on the left”.  Award 0.7 points for this.
4. Return a progressive score (max 1.0).  Print detailed diagnostics and the final
   reward in the required format:  "REWARD: X.X".

Notes:
• The concept of “Anchor to Paragraph” isn’t explicitly represented in PPTX like
  it is in Word; for this Impress task we equate the requirement to the shape
  being placed independently (i.e., it appears as a picture shape) and located
  flush left.  Verifying that spatial positioning achieved the instruction is the
  measurable part here.
• No points are ever given just for being able to open the file or for natural
  conditions (e.g., slide existence).
"""

def verify_anchor_second_image(file_path: str) -> float:
    print(f"Loading presentation: {file_path}")

    # ------------------------------------------------------------------
    # 0. Prerequisite – file existence & readability (no score impact)
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File not found – task failed")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Unable to open presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    total_score = 0.0  # progressive scoring

    # ------------------------------------------------------------------
    # 1. Gather all picture shapes across the presentation
    # ------------------------------------------------------------------
    pictures = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pictures.append({
                    'slide_idx': slide_idx,
                    'shape': shape
                })

    print(f"Total picture shapes found: {len(pictures)}")

    # ------------------------------------------------------------------
    # 2. Requirement A – at least two images present (0.3 points)
    # ------------------------------------------------------------------
    if len(pictures) >= 2:
        total_score += 0.3
        print("✓ Found at least two images (0.3)")
    else:
        print("✗ Less than two images – cannot meet instruction")
        print(f"REWARD: {total_score}")
        return total_score

    # ------------------------------------------------------------------
    # 3. Requirement B – second image positioned on the left (0.7 points)
    # ------------------------------------------------------------------
    second_pic_shape = pictures[1]['shape']  # 0-based index → second in order
    left_emu = second_pic_shape.left
    slide_width_emu = prs.slide_width
    threshold_emu = int(0.05 * slide_width_emu)  # 5 % of slide width

    print(f"Second image left position: {left_emu} EMU (threshold ≤ {threshold_emu})")

    if left_emu <= threshold_emu:
        total_score += 0.7
        print("✓ Second image is positioned near the left edge (0.7)")
    else:
        print("✗ Second image is NOT positioned on the left – no points for positioning")

    # ------------------------------------------------------------------
    # 4. Final score (capped at 1.0)
    # ------------------------------------------------------------------
    final_score = min(total_score, 1.0)
    print(f"REWARD: {final_score}")
    return final_score

# ----------------------------------------------------------------------
# Execute verification when the script is run directly
# ----------------------------------------------------------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/anchor_the_second_image_to_paragraph_and_position_it_on_the_left.pptx"
    verify_anchor_second_image(FILE_PATH)
