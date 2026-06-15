"""
Reward Script: Create a 4-slide team presentation with blue theme
Task ID: impress_wf_002
Domain: libreoffice_impress
Scoring:
  Component 1 (0.15): Exactly 4 slides
  Component 2 (0.25): Slide 1 title='Meet Our Team', subtitle='Engineering Department'
  Component 3 (0.30): Slides 2-4 each have a title and 3 bullet points
  Component 4 (0.30): All slides have background #2C5F8A
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_002'


def get_slide_background_rgb(slide):
    """Get background color as uppercase hex string, or None."""
    fill = slide.background.fill
    if fill.type is not None and fill.type == 1:  # SOLID
        return str(fill.fore_color.rgb).upper()
    elif fill.type == 5:  # inherited from master
        master_fill = slide.slide_layout.slide_master.background.fill
        if master_fill.type == 1:
            return str(master_fill.fore_color.rgb).upper()
    return None


def get_all_text_shapes(slide):
    """Recursively get all shapes with text frames."""
    def extract(shape):
        results = []
        if hasattr(shape, "text") and hasattr(shape, "text_frame"):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.extend(extract(sub))
        return results
    out = []
    for shape in slide.shapes:
        out.extend(extract(shape))
    return out


def get_nonempty_paragraphs(shape):
    """Return paragraphs with non-empty text from a shape's text frame."""
    if not shape.has_text_frame:
        return []
    return [p for p in shape.text_frame.paragraphs if p.text.strip()]


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Exactly 4 slides (0.15 points)
    try:
        num_slides = len(prs.slides)
        if num_slides == 4:
            print(f"PASS: Component 1 — Slide count is 4 (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — Expected 4 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 1 has title 'Meet Our Team' and subtitle 'Engineering Department' (0.25 points)
    try:
        if len(prs.slides) >= 1:
            slide1 = prs.slides[0]
            shapes = get_all_text_shapes(slide1)
            # Collect all text from shapes
            all_texts = []
            for shape in shapes:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        all_texts.append(t)

            has_title = any('meet our team' in t.lower() for t in all_texts)
            has_subtitle = any('engineering department' in t.lower() for t in all_texts)

            if has_title and has_subtitle:
                print(f"PASS: Component 2 — Slide 1 has title 'Meet Our Team' and subtitle 'Engineering Department' (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 2 — Slide 1 texts: {all_texts}. has_title={has_title}, has_subtitle={has_subtitle}")
        else:
            print(f"FAIL: Component 2 — No slides available")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slides 2-4 each have a title and 3 bullet points (0.30 points)
    # Award 0.10 per slide that meets criteria
    try:
        comp3_score = 0.0
        if len(prs.slides) >= 4:
            for slide_idx in range(1, 4):
                slide = prs.slides[slide_idx]
                shapes = get_all_text_shapes(slide)

                # Find title shape (first shape with text, typically a placeholder title)
                title_text = None
                bullet_texts = []

                for shape in shapes:
                    nonempty = get_nonempty_paragraphs(shape)
                    if not nonempty:
                        continue
                    # First shape with text is treated as title candidate
                    # Subsequent shape or multi-paragraph shape provides bullets
                    if title_text is None and len(nonempty) == 1:
                        title_text = nonempty[0].text.strip()
                    elif title_text is None and len(nonempty) > 1:
                        # First para could be title, rest are bullets
                        title_text = nonempty[0].text.strip()
                        bullet_texts.extend([p.text.strip() for p in nonempty[1:]])
                    else:
                        # This shape provides bullet points
                        bullet_texts.extend([p.text.strip() for p in nonempty])

                has_title = title_text is not None and len(title_text) > 0
                has_3_bullets = len(bullet_texts) >= 3

                if has_title and has_3_bullets:
                    print(f"PASS: Component 3.{slide_idx} — Slide {slide_idx+1} has title '{title_text}' and {len(bullet_texts)} bullet points (0.10 pts)")
                    comp3_score += 0.10
                else:
                    print(f"FAIL: Component 3.{slide_idx} — Slide {slide_idx+1}: title={repr(title_text)}, bullets={len(bullet_texts)}")

        else:
            print(f"FAIL: Component 3 — Need at least 4 slides")

        if comp3_score > 0:
            total_score += comp3_score
            print(f"  Component 3 subtotal: {comp3_score:.2f}/0.30")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: All slides have background color #2C5F8A (0.30 points)
    try:
        slides_with_bg = 0
        num_slides_to_check = min(len(prs.slides), 4)
        if num_slides_to_check == 0:
            print(f"FAIL: Component 4 — No slides to check")
        else:
            for i in range(num_slides_to_check):
                slide = prs.slides[i]
                bg_color = get_slide_background_rgb(slide)
                if bg_color == '2C5F8A':
                    slides_with_bg += 1
                    print(f"  Slide {i+1} background: #{bg_color} — correct")
                else:
                    print(f"  Slide {i+1} background: {bg_color} — expected #2C5F8A")

            if slides_with_bg == num_slides_to_check and num_slides_to_check == 4:
                print(f"PASS: Component 4 — All 4 slides have background #2C5F8A (0.30 pts)")
                total_score += 0.30
            else:
                print(f"FAIL: Component 4 — {slides_with_bg}/{num_slides_to_check} slides have correct background")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/Desktop/Team_Intro.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
