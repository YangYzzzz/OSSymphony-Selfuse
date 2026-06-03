"""
Initial Setup: Create a 5-slide Impress presentation with build-up animations
Task ID: impress_el_039
Domain: libreoffice_impress

Slide 2 has 3 bullet points that appear one by one (3 animation steps).
Slide 4 has 2 items that appear one by one (2 animation steps).
Other slides have no animations.
"""

import os
import shlex
import subprocess
import time
import zipfile
import shutil
import tempfile
from xml.etree import ElementTree as ET

# python-pptx for base structure
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_el_039'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Quarterly Business Review"
    slide1.placeholders[1].text = "FY2025 Q1 Performance Summary\nPrepared by Strategic Planning Team"

    # ---- Slide 2: Three bullet points with build-up animations ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Key Achievements"
    body2 = slide2.placeholders[1]
    tf2 = body2.text_frame
    tf2.clear()

    p1 = tf2.paragraphs[0]
    p1.text = "Revenue grew 23% year-over-year to $4.2M in Q1"
    p1.font.size = Pt(20)

    p2 = tf2.add_paragraph()
    p2.text = "Customer acquisition cost decreased by 15% through improved targeting"
    p2.font.size = Pt(20)

    p3 = tf2.add_paragraph()
    p3.text = "Net Promoter Score reached an all-time high of 72"
    p3.font.size = Pt(20)

    # ---- Slide 3: Regular content (no animations) ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Market Analysis"
    body3 = slide3.placeholders[1]
    tf3 = body3.text_frame
    tf3.clear()
    lines3 = [
        "Total addressable market expanded to $12B globally",
        "Primary competitors lost 3.2% market share combined",
        "Emerging markets showed strongest growth at 31% YoY",
        "Digital channel adoption accelerated across all segments",
    ]
    for i, line in enumerate(lines3):
        if i == 0:
            tf3.paragraphs[0].text = line
            tf3.paragraphs[0].font.size = Pt(20)
        else:
            p = tf3.add_paragraph()
            p.text = line
            p.font.size = Pt(20)

    # ---- Slide 4: Two items with build-up animations ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Strategic Priorities"
    body4 = slide4.placeholders[1]
    tf4 = body4.text_frame
    tf4.clear()

    p4a = tf4.paragraphs[0]
    p4a.text = "Expand enterprise sales team by 40% in APAC region"
    p4a.font.size = Pt(20)

    p4b = tf4.add_paragraph()
    p4b.text = "Launch self-service analytics platform for mid-market clients"
    p4b.font.size = Pt(20)

    # ---- Slide 5: Summary/Closing ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Next Steps & Timeline"
    body5 = slide5.placeholders[1]
    tf5 = body5.text_frame
    tf5.clear()
    lines5 = [
        "Board review scheduled for April 15, 2025",
        "Budget reallocation proposals due by April 30",
        "Q2 kickoff all-hands meeting on May 5",
    ]
    for i, line in enumerate(lines5):
        if i == 0:
            tf5.paragraphs[0].text = line
            tf5.paragraphs[0].font.size = Pt(20)
        else:
            p = tf5.add_paragraph()
            p.text = line
            p.font.size = Pt(20)

    # Save initial pptx
    prs.save(OUTPUT)
    print(f'Base presentation saved: {OUTPUT}')

    # Now inject animations via XML manipulation into the saved pptx
    inject_animations(OUTPUT)
    print(f'Animations injected into: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=3.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


def inject_animations(pptx_path):
    """Inject appear-one-by-one animations into slide 2 (3 steps) and slide 4 (2 steps)."""
    # We need to add animation XML to the slide XML files inside the pptx zip

    ns = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }

    # Register namespaces to preserve them
    for prefix, uri in ns.items():
        ET.register_namespace(prefix, uri)
    # Also register common namespaces
    ET.register_namespace('p14', 'http://schemas.microsoft.com/office/powerpoint/2010/main')
    ET.register_namespace('mc', 'http://schemas.openxmlformats.org/markup-compatibility/2006')
    ET.register_namespace('p15', 'http://schemas.microsoft.com/office/powerpoint/2012/main')

    tmp_dir = tempfile.mkdtemp()
    tmp_pptx = os.path.join(tmp_dir, 'modified.pptx')

    # Extract, modify, repack
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        with zipfile.ZipFile(tmp_pptx, 'w') as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)

                if item.filename == 'ppt/slides/slide2.xml':
                    data = add_appear_animations(data, ns, num_paragraphs=3, shape_idx=1)
                elif item.filename == 'ppt/slides/slide4.xml':
                    data = add_appear_animations(data, ns, num_paragraphs=2, shape_idx=1)

                zout.writestr(item, data)

    shutil.move(tmp_pptx, pptx_path)
    shutil.rmtree(tmp_dir, ignore_errors=True)


def add_appear_animations(slide_xml_bytes, ns, num_paragraphs, shape_idx):
    """Add 'Appear' animation to each paragraph in a content placeholder.

    This creates click-triggered appear animations for each paragraph,
    making them build up one by one.
    """
    root = ET.fromstring(slide_xml_bytes)

    # Find the content shape (second shape = placeholders[1])
    sp_tree = root.find('.//p:cSld/p:spTree', ns)
    shapes = sp_tree.findall('p:sp', ns)

    if shape_idx >= len(shapes):
        return slide_xml_bytes

    target_shape = shapes[shape_idx]
    # Get the shape's spId
    nvSpPr = target_shape.find('p:nvSpPr', ns)
    cNvPr = nvSpPr.find('p:cNvPr', ns)
    shape_id = cNvPr.get('id')

    # Build the animation timing XML
    # Each paragraph gets an "appear" effect on click
    p_ns = ns['p']
    a_ns = ns['a']

    # Create timing node
    timing = ET.SubElement(root, f'{{{p_ns}}}timing')
    tnLst = ET.SubElement(timing, f'{{{p_ns}}}tnLst')

    # Main sequence par
    par_main = ET.SubElement(tnLst, f'{{{p_ns}}}par')
    cTn_main = ET.SubElement(par_main, f'{{{p_ns}}}cTn')
    cTn_main.set('id', '1')
    cTn_main.set('dur', 'indefinite')
    cTn_main.set('restart', 'never')
    cTn_main.set('nodeType', 'tmRoot')

    childTnLst_main = ET.SubElement(cTn_main, f'{{{p_ns}}}childTnLst')

    # Sequence container
    seq = ET.SubElement(childTnLst_main, f'{{{p_ns}}}seq')
    seq.set('concurrent', '1')
    seq.set('nextAc', 'seek')

    cTn_seq = ET.SubElement(seq, f'{{{p_ns}}}cTn')
    cTn_seq.set('id', '2')
    cTn_seq.set('dur', 'indefinite')
    cTn_seq.set('nodeType', 'mainSeq')

    childTnLst_seq = ET.SubElement(cTn_seq, f'{{{p_ns}}}childTnLst')

    next_id = 3
    for para_idx in range(num_paragraphs):
        # Each paragraph animation is a par with click trigger
        par_anim = ET.SubElement(childTnLst_seq, f'{{{p_ns}}}par')
        cTn_par = ET.SubElement(par_anim, f'{{{p_ns}}}cTn')
        cTn_par.set('id', str(next_id))
        next_id += 1
        cTn_par.set('fill', 'hold')

        stCondLst = ET.SubElement(cTn_par, f'{{{p_ns}}}stCondLst')
        cond = ET.SubElement(stCondLst, f'{{{p_ns}}}cond')
        cond.set('delay', '0')

        childTnLst_par = ET.SubElement(cTn_par, f'{{{p_ns}}}childTnLst')

        # Inner par for the effect
        par_inner = ET.SubElement(childTnLst_par, f'{{{p_ns}}}par')
        cTn_inner = ET.SubElement(par_inner, f'{{{p_ns}}}cTn')
        cTn_inner.set('id', str(next_id))
        next_id += 1
        cTn_inner.set('presetID', '1')  # 1 = Appear
        cTn_inner.set('presetClass', 'entr')
        cTn_inner.set('presetSubtype', '0')
        cTn_inner.set('fill', 'hold')
        cTn_inner.set('nodeType', 'clickEffect')

        stCondLst2 = ET.SubElement(cTn_inner, f'{{{p_ns}}}stCondLst')
        cond2 = ET.SubElement(stCondLst2, f'{{{p_ns}}}cond')
        cond2.set('delay', '0')

        childTnLst_inner = ET.SubElement(cTn_inner, f'{{{p_ns}}}childTnLst')

        # Set effect - make paragraph visible
        set_elem = ET.SubElement(childTnLst_inner, f'{{{p_ns}}}set')
        cBhvr = ET.SubElement(set_elem, f'{{{p_ns}}}cBhvr')
        cTn_set = ET.SubElement(cBhvr, f'{{{p_ns}}}cTn')
        cTn_set.set('id', str(next_id))
        next_id += 1
        cTn_set.set('dur', '1')
        cTn_set.set('fill', 'hold')

        stCondLst3 = ET.SubElement(cTn_set, f'{{{p_ns}}}stCondLst')
        cond3 = ET.SubElement(stCondLst3, f'{{{p_ns}}}cond')
        cond3.set('delay', '0')

        tgtEl = ET.SubElement(cBhvr, f'{{{p_ns}}}tgtEl')
        spTgt = ET.SubElement(tgtEl, f'{{{p_ns}}}spTgt')
        spTgt.set('spid', shape_id)
        txEl = ET.SubElement(spTgt, f'{{{p_ns}}}txEl')
        pRg = ET.SubElement(txEl, f'{{{p_ns}}}pRg')
        pRg.set('st', str(para_idx))
        pRg.set('end', str(para_idx))

        attrNameLst = ET.SubElement(cBhvr, f'{{{p_ns}}}attrNameLst')
        attrName = ET.SubElement(attrNameLst, f'{{{p_ns}}}attrName')
        attrName.text = 'style.visibility'

        to_elem = ET.SubElement(set_elem, f'{{{p_ns}}}to')
        val = ET.SubElement(to_elem, f'{{{p_ns}}}strVal')
        val.set('val', 'visible')

    # prevCondLst and nextCondLst for the sequence
    prevCondLst = ET.SubElement(seq, f'{{{p_ns}}}prevCondLst')
    prev_cond = ET.SubElement(prevCondLst, f'{{{p_ns}}}cond')
    prev_cond.set('evt', 'onPrev')
    prev_cond.set('delay', '0')
    prev_tgt = ET.SubElement(prev_cond, f'{{{p_ns}}}tgtEl')
    ET.SubElement(prev_tgt, f'{{{p_ns}}}sldTgt')

    nextCondLst = ET.SubElement(seq, f'{{{p_ns}}}nextCondLst')
    next_cond = ET.SubElement(nextCondLst, f'{{{p_ns}}}cond')
    next_cond.set('evt', 'onNext')
    next_cond.set('delay', '0')
    next_tgt = ET.SubElement(next_cond, f'{{{p_ns}}}tgtEl')
    ET.SubElement(next_tgt, f'{{{p_ns}}}sldTgt')

    return ET.tostring(root, xml_declaration=True, encoding='UTF-8')


create_initial()
