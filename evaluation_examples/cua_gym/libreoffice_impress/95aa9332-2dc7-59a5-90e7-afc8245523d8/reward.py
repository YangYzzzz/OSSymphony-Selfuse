"""
FINAL REWARD SCRIPT - SUCCESS
Task: In LibreOffice Impress I spotted that the title on slide 281 isn’t using the right style. Could you walk me through how to switch that one title to Liberation Sans Narrow, 46 pt, bold?
Generated: 2025-09-10 21:54:50
Status: success
Model: azure-o3
Total Steps: 10
"""

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Pt
import os


def verify_slide_281_title_style(file_path: str) -> float:
    """Verify that the title on slide 281 is Liberation Sans Narrow, 46 pt, bold.

    Returns a progressive score between 0.0 and 1.0.
    A perfect match on all three attributes (font name, size, bold) yields 1.0.
    """
    print(f"Starting verification for: {file_path}\n")
    total_score: float = 0.0
    max_score: float = 1.0

    # ------------------------------------------------------------------
    # 1. Load presentation (NO POINTS: prerequisite, but fail fast)
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File not found.")
        return 0.0

    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded successfully with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        return 0.0

    # ------------------------------------------------------------------
    # 2. Ensure slide 281 exists (index 280) – task-specific, earns points
    # ------------------------------------------------------------------
    slide_idx = 280  # 0-based index
    if slide_idx >= len(prs.slides):
        print("✗ Slide 281 does not exist in the presentation")
        return 0.0

    slide = prs.slides[slide_idx]
    print("✓ Slide 281 exists")
    total_score += 0.2  # 20 % for target slide presence (task-specific)

    # ------------------------------------------------------------------
    # 3. Locate the title placeholder on slide 281 – earns points
    # ------------------------------------------------------------------
    title_shape = None
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                title_shape = shape
                break

    if title_shape is None:
        print("✗ Could not find a title placeholder on slide 281")
        return total_score  # cannot verify style further

    print("✓ Title placeholder located")
    total_score += 0.2  # 20 % for correctly identifying the title shape

    # ------------------------------------------------------------------
    # 4. Verify font attributes across all non-empty runs
    # ------------------------------------------------------------------
    expected_font_name = "liberation sans narrow"
    expected_size_pt = 46
    expected_size_emu = int(Pt(expected_size_pt))  # size in EMU units
    expected_bold = True

    name_ok = True
    size_ok = True
    bold_ok = True
    found_runs = False

    for para in title_shape.text_frame.paragraphs:
        for run in para.runs:
            if not run.text.strip():
                continue  # skip empty runs
            found_runs = True
            font = run.font

            # -------- Font name check --------
            if font.name is None or font.name.lower() != expected_font_name:
                name_ok = False

            # -------- Font size check (tolerance 1 pt) --------
            if font.size is None or abs(int(font.size) - expected_size_emu) > 12700:  # 1 pt ≈ 12 700 EMU
                size_ok = False

            # -------- Bold check --------
            if font.bold is not True:
                bold_ok = False

    if not found_runs:
        print("✗ No text runs found inside the title placeholder – cannot verify style")
        return total_score

    # ------------------------------------------------------------------
    # 5. Progressive scoring for each correctly set attribute
    # ------------------------------------------------------------------
    if name_ok:
        print("✓ Font name is Liberation Sans Narrow")
        total_score += 0.34
    else:
        print("✗ Font name is not Liberation Sans Narrow")

    if size_ok:
        print("✓ Font size is 46 pt")
        total_score += 0.33
    else:
        print("✗ Font size is not 46 pt")

    if bold_ok:
        print("✓ Font is bold as required")
        total_score += 0.33
    else:
        print("✗ Font is not bold")

    # Clamp to [0, 1]
    final_score = min(total_score, max_score)

    print(f"\nFINAL SCORE: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score


if __name__ == "__main__":
    FILE_PATH = "/home/user/in_libreoffice_impress_i_spotted_that_the_title_on_slide_281_isnt_using_the_right_style_could_you_wa_golden.pptx"
    verify_slide_281_title_style(FILE_PATH)

