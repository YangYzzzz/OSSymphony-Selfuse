"""
FINAL REWARD SCRIPT - SUCCESS
Task: Hey, I'm working on slide 8 of my presentation and I've got three images that I'd like to line up perfectly along the top edge. What's the best way to do this in LibreOffice Impress?
Generated: 2025-08-07 08:59:48
Status: success
Model: o4-mini
Total Steps: 2
"""

import os
from pptx import Presentation

def verify_alignment_on_slide8(file_path):
    print("Checking task: Align three images on slide 8 along the top edge")
    total_score = 0.0
    max_score = 1.0

    # 1. Check file existence (0.2)
    try:
        if os.path.exists(file_path):
            print("✓ File exists (0.2 points)")
            total_score += 0.2
        else:
            print(f"✗ File not found: {file_path} (0 points)")
            print(f"REWARD: {total_score:.1f}")
            return total_score
    except Exception as e:
        print(f"✗ Error checking file existence: {e}")
        print(f"REWARD: {total_score:.1f}")
        return total_score

    # 2. Load presentation (0.1)
    try:
        prs = Presentation(file_path)
        print(f"✓ Successfully loaded presentation with {len(prs.slides)} slides (0.1 points)")
        total_score += 0.1
    except Exception as e:
        print(f"✗ Error loading presentation: {e} (0 points)")
        print(f"REWARD: {total_score:.1f}")
        return total_score

    # 3. Check slide count >= 8 (0.1)
    slide_count = len(prs.slides)
    if slide_count >= 8:
        print(f"✓ Presentation has at least 8 slides: {slide_count} (0.1 points)")
        total_score += 0.1
    else:
        print(f"✗ Not enough slides: Found {slide_count} slides (0 points)")
        print(f"REWARD: {total_score:.1f}")
        return total_score

    # 4. Verify exactly 3 images on slide 8 (0.2)
    slide = prs.slides[7]
    pics = [shape for shape in slide.shapes if hasattr(shape, 'image')]
    num_pics = len(pics)
    if num_pics == 3:
        print("✓ Exactly 3 images found on slide 8 (0.2 points)")
        total_score += 0.2
    else:
        print(f"✗ Expected 3 images, found {num_pics} (0 points)")
        print(f"REWARD: {total_score:.1f}")
        return total_score

    # 5. Verify all images share the same top position (0.4)
    tops = [shape.top for shape in pics]
    unique_tops = set(tops)
    if len(unique_tops) == 1:
        print(f"✓ All images share the same top position ({tops[0]}) (0.4 points)")
        total_score += 0.4
    else:
        print(f"✗ Images have differing top positions: {tops} (0 points)")
        print(f"REWARD: {total_score:.1f}")
        return total_score

    # Final score
    final_score = min(total_score, max_score)
    print(f"Total score: {total_score:.1f}/{max_score}")
    print(f"REWARD: {final_score:.1f}")
    return final_score

if __name__ == '__main__':
    file_path = '/home/user/hey_im_working_on_slide_8_of_my_presentation_and_ive_got_three_images_that_id_like_to_line_up_perfec.pptx'
    verify_alignment_on_slide8(file_path)
