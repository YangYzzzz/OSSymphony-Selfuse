"""
Reward Script: Slide master template with dark blue top bar, gold line, department text, and date placeholder
Task ID: impress_teach_056
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Dark blue rectangle bar across top of slide master (full width, 0.75in tall, #0D47A1)
  Component 2 (0.20): Gold line (#FFD700) below the top bar
  Component 3 (0.30): White 12pt text 'Department of Computer Science' in the top bar area
  Component 4 (0.20): Date placeholder in the bottom-left corner of the slide master
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_056'
FILE_PATH = os.path.join(WORKDIR, f'{TASK_ID}.pptx')

# Constants for verification
SLIDE_WIDTH = 9144000  # 10 inches in EMU
FULL_WIDTH_TOLERANCE = 0.05  # 5% tolerance for "full width"
HEIGHT_075IN = Inches(0.75)  # 685800 EMU
HEIGHT_TOLERANCE = 0.10  # 10% tolerance for height
COLOR_DARK_BLUE = '0D47A1'
COLOR_GOLD = 'FFD700'
COLOR_WHITE = 'FFFFFF'
FONT_SIZE_12PT = Pt(12)  # 152400 EMU


def get_shape_fill_color_xml(shape):
    """Get solid fill color from shape XML."""
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is not None:
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is not None:
            srgbClr = solidFill.find(qn('a:srgbClr'))
            if srgbClr is not None:
                return srgbClr.get('val')
    return None


def get_shape_line_color_xml(shape):
    """Get line/outline color from shape XML."""
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is not None:
        ln = spPr.find(qn('a:ln'))
        if ln is not None:
            lnFill = ln.find(qn('a:solidFill'))
            if lnFill is not None:
                lnClr = lnFill.find(qn('a:srgbClr'))
                if lnClr is not None:
                    return lnClr.get('val')
    return None


def get_shape_preset_geom(shape):
    """Get preset geometry type from shape XML."""
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is not None:
        pg = spPr.find(qn('a:prstGeom'))
        if pg is not None:
            return pg.get('prst')
    return None


def is_approx(val, expected, tolerance=0.10):
    """Check if val is approximately equal to expected within relative tolerance."""
    if expected == 0:
        return val == 0
    return abs(val - expected) / max(abs(val), abs(expected)) <= tolerance


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Get the first slide master
    if len(prs.slide_masters) == 0:
        print("CRITICAL: No slide masters found")
        print("REWARD: 0.0")
        return 0.0

    master = prs.slide_masters[0]
    slide_width = prs.slide_width

    # Component 1: Dark blue rectangle bar across top (0.30 points)
    # Must be: full width, ~0.75in tall, fill color #0D47A1, positioned at top (top=0)
    try:
        found_topbar = False
        for shape in master.shapes:
            geom = get_shape_preset_geom(shape)
            fill_color = get_shape_fill_color_xml(shape)

            # Look for a rectangle with dark blue fill at the top
            if (geom == 'rect' and
                fill_color is not None and
                fill_color.upper() == COLOR_DARK_BLUE):

                # Check position: should be at top (top ~= 0)
                if shape.top <= Inches(0.1):
                    # Check width: should be full slide width
                    width_ratio = shape.width / slide_width
                    if width_ratio >= 0.95:
                        # Check height: should be ~0.75 inches
                        if is_approx(shape.height, HEIGHT_075IN, tolerance=0.15):
                            found_topbar = True
                            print(f"PASS: Component 1 — Dark blue top bar found: "
                                  f"fill={fill_color}, width={shape.width} "
                                  f"({width_ratio:.2%} of slide), "
                                  f"height={shape.height} EMU "
                                  f"(~{shape.height/914400:.2f}in) (0.30 pts)")
                            total_score += 0.30
                            break

        if not found_topbar:
            print(f"FAIL: Component 1 — No dark blue (#0D47A1) rectangle bar found "
                  f"at top of slide master (full width, ~0.75in tall)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Gold line below the top bar (0.20 points)
    # Must be: line shape with #FFD700 color, positioned at/near bottom of the top bar
    try:
        found_gold_line = False
        for shape in master.shapes:
            geom = get_shape_preset_geom(shape)
            line_color = get_shape_line_color_xml(shape)

            # Look for a line with gold color
            if (geom == 'line' and
                line_color is not None and
                line_color.upper() == COLOR_GOLD):
                # Check that it's positioned near the bottom of the top bar area
                # The top bar is ~0.75in (685800 EMU), so the line should be around there
                if shape.top >= Inches(0.5) and shape.top <= Inches(1.2):
                    # Check width is roughly full width
                    width_ratio = shape.width / slide_width
                    if width_ratio >= 0.90:
                        found_gold_line = True
                        print(f"PASS: Component 2 — Gold line found: "
                              f"color={line_color}, top={shape.top} EMU "
                              f"(~{shape.top/914400:.2f}in), "
                              f"width={shape.width} ({width_ratio:.2%} of slide) (0.20 pts)")
                        total_score += 0.20
                        break

        if not found_gold_line:
            # Also check for a very thin rectangle with gold fill as alternative
            for shape in master.shapes:
                geom = get_shape_preset_geom(shape)
                fill_color = get_shape_fill_color_xml(shape)
                line_color = get_shape_line_color_xml(shape)
                color = fill_color or line_color

                if color is not None and color.upper() == COLOR_GOLD:
                    if shape.top >= Inches(0.5) and shape.top <= Inches(1.2):
                        # Thin rectangle or any shape acting as a line
                        if shape.height <= Inches(0.15):
                            width_ratio = shape.width / slide_width
                            if width_ratio >= 0.90:
                                found_gold_line = True
                                print(f"PASS: Component 2 — Gold line/thin shape found: "
                                      f"color={color}, height={shape.height} EMU, "
                                      f"top={shape.top} EMU (0.20 pts)")
                                total_score += 0.20
                                break

        if not found_gold_line:
            print(f"FAIL: Component 2 — No gold (#FFD700) line found below the top bar")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: White 12pt text 'Department of Computer Science' in the top bar (0.30 points)
    # Must be: text containing the department name, white color, ~12pt size, positioned within top bar area
    try:
        found_dept_text = False
        for shape in master.shapes:
            if not (hasattr(shape, 'has_text_frame') and shape.has_text_frame):
                continue
            # Shape should be in the top bar area (top < 1 inch)
            if shape.top > Inches(1.0):
                continue

            full_text = ''
            for para in shape.text_frame.paragraphs:
                full_text += para.text

            if 'Department of Computer Science' not in full_text:
                continue

            # Found the text - now check font properties
            text_ok = False
            color_ok = False
            size_ok = False

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if 'Department' not in run.text and 'Computer Science' not in run.text:
                        continue
                    # Check color
                    try:
                        if run.font.color.type is not None:
                            rgb = str(run.font.color.rgb).upper()
                            if rgb == COLOR_WHITE:
                                color_ok = True
                    except Exception:
                        pass
                    # Check size (~12pt = 152400 EMU)
                    if run.font.size is not None:
                        if is_approx(run.font.size, FONT_SIZE_12PT, tolerance=0.15):
                            size_ok = True

            text_ok = True  # We already confirmed text contains the department name

            sub_score = 0.0
            if text_ok:
                sub_score += 0.10
            if color_ok:
                sub_score += 0.10
            if size_ok:
                sub_score += 0.10

            if sub_score > 0:
                found_dept_text = True
                print(f"PASS: Component 3 — Department text found: "
                      f"text={'OK' if text_ok else 'FAIL'}, "
                      f"white={'OK' if color_ok else 'FAIL'}, "
                      f"12pt={'OK' if size_ok else 'FAIL'} ({sub_score} pts)")
                total_score += sub_score
                break

        if not found_dept_text:
            print(f"FAIL: Component 3 — 'Department of Computer Science' not found "
                  f"in white 12pt text on slide master top bar area")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Date placeholder in bottom-left corner (0.20 points)
    # Must be: a date-type placeholder positioned in the bottom-left area of the slide
    # The task says "date placeholder in the bottom-left corner"
    # We check for a placeholder with date content or date type in the lower portion
    try:
        found_date_ph = False
        slide_height = prs.slide_height

        # Count non-standard shapes on master (shapes beyond the default 5 placeholders)
        # We need a date placeholder specifically in bottom-left
        for shape in master.shapes:
            # Check if it's a placeholder
            is_placeholder = (str(shape.shape_type) == 'PLACEHOLDER (14)')

            if not is_placeholder:
                continue

            # Must be in bottom area (bottom 20% of slide)
            if shape.top < slide_height * 0.75:
                continue

            # Must be in left area (left 50% of slide)
            if shape.left > slide_width * 0.50:
                continue

            # Check if this is a date-related placeholder
            # The default Date Placeholder 3 exists on initial too (idx=2),
            # so we need to check if there's a NEW/different date placeholder
            # or if the existing one was moved
            ph_format = shape.placeholder_format
            ph_idx = ph_format.idx if ph_format else None
            ph_type = ph_format.type if ph_format else None

            # Look for date type (16) or a placeholder with date-like name/content
            # that is NOT one of the original 5 (idx 0,1,2,3,4)
            if ph_idx is not None and ph_idx not in (0, 1, 2, 3, 4):
                # This is a new placeholder added by the task
                # Check if it has date content or is date-typed
                has_date_content = False
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    # Date patterns
                    import re
                    if re.search(r'\d{1,2}/\d{1,2}/\d{2,4}', text) or \
                       re.search(r'\d{4}-\d{2}-\d{2}', text) or \
                       'date' in shape.name.lower():
                        has_date_content = True

                if has_date_content or 'date' in shape.name.lower():
                    found_date_ph = True
                    print(f"PASS: Component 4 — New date placeholder at bottom-left: "
                          f"name={shape.name}, idx={ph_idx}, "
                          f"left={shape.left}, top={shape.top} (0.20 pts)")
                    total_score += 0.20
                    break

        if not found_date_ph:
            # Fallback: check if any placeholder in bottom-left has date content
            # and is different from what the initial had (new name or new position)
            for shape in master.shapes:
                is_placeholder = (str(shape.shape_type) == 'PLACEHOLDER (14)')
                if not is_placeholder:
                    continue
                if shape.top < slide_height * 0.75:
                    continue
                if shape.left > slide_width * 0.50:
                    continue

                # Check for a placeholder named differently from defaults
                default_names = {'Title Placeholder 1', 'Text Placeholder 2',
                                 'Date Placeholder 3', 'Footer Placeholder 4',
                                 'Slide Number Placeholder 5'}
                if shape.name not in default_names:
                    if 'date' in shape.name.lower() or 'Date' in shape.name:
                        found_date_ph = True
                        print(f"PASS: Component 4 — Date placeholder at bottom-left: "
                              f"name={shape.name}, left={shape.left}, top={shape.top} (0.20 pts)")
                        total_score += 0.20
                        break

        if not found_date_ph:
            print(f"FAIL: Component 4 — No new date placeholder found in bottom-left corner of slide master")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice
def persist_app_state(domain):
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
persist_app_state("libreoffice_impress")

if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(FILE_PATH)
