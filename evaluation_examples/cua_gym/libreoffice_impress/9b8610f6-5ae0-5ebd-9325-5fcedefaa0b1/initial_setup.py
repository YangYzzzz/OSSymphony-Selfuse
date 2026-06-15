"""
Initial Setup: Create a 5-slide presentation with a plain white master.
Task ID: impress_teach_056
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_056'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "CS Department Overview"
    slide1.placeholders[1].text = "Academic Year 2025-2026"

    # --- Slide 2: Faculty ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Faculty Members"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Dr. Elena Vasquez - Machine Learning"
    body2.add_paragraph().text = "Dr. Rajesh Patel - Distributed Systems"
    body2.add_paragraph().text = "Dr. Mei-Lin Chang - Computer Vision"
    body2.add_paragraph().text = "Dr. Samuel Okonkwo - Cybersecurity"
    body2.add_paragraph().text = "Dr. Anna Kowalski - Human-Computer Interaction"

    # --- Slide 3: Research Areas ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Research Focus Areas"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Artificial Intelligence & Deep Learning"
    body3.add_paragraph().text = "Cloud Computing & Edge Infrastructure"
    body3.add_paragraph().text = "Quantum Computing Algorithms"
    body3.add_paragraph().text = "Software Engineering & DevOps"
    body3.add_paragraph().text = "Bioinformatics & Computational Biology"

    # --- Slide 4: Student Statistics ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Student Enrollment 2025"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Undergraduate Students: 342"
    body4.add_paragraph().text = "Master's Students: 128"
    body4.add_paragraph().text = "PhD Candidates: 47"
    body4.add_paragraph().text = "Post-Doctoral Researchers: 12"
    body4.add_paragraph().text = "Total Faculty-to-Student Ratio: 1:18"

    # --- Slide 5: Upcoming Events ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Key Events This Semester"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "March 15 - Annual Research Symposium"
    body5.add_paragraph().text = "April 2 - Industry Partnership Day"
    body5.add_paragraph().text = "April 20 - Graduate Student Showcase"
    body5.add_paragraph().text = "May 8 - Commencement Ceremony"
    body5.add_paragraph().text = "May 22 - Summer Research Kickoff"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
