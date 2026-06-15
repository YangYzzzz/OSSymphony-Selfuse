"""
Reward Script: Edit slide 2 title from 'Scope' to 'Project Scope and Objectives'
Task ID: impress_text_039
Domain: libreoffice_impress
Scoring:
  Component 1 (0.6): Slide 2 title placeholder text is exactly 'Project Scope and Objectives'
  Component 2 (0.4): Slide 2 title is correct AND body text is unchanged AND slide count is 5
                     (compound check - ensures no collateral damage from the edit)
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'impress_text_039'

# Expected values derived from task_config context
EXPECTED_SLIDE2_TITLE = 'Project Scope and Objectives'
EXPECTED_SLIDE_COUNT = 5
EXPECTED_SLIDE2_BODY_LINES = [
    'Define the boundaries of the infrastructure modernization effort',
    'Covers data center consolidation across 3 regional offices',
    'Network redesign for improved redundancy and throughput',
    'Cloud migration of 40 legacy on-premise workloads',
    'Excludes customer-facing application redesign',
]


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: file must be loadable
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition gate: presentation must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"CRITICAL: Expected at least 2 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Helper: get title from slide by placeholder index 0
    def get_slide_title(slide):
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                return shape.text_frame.text.strip()
        return None

    # Helper: get body text lines from slide placeholder index 1
    def get_slide_body_lines(slide):
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 1:
                raw = shape.text_frame.text.strip()
                return [line.strip() for line in raw.split('\n') if line.strip()]
        return []

    # Component 1: Slide 2 title is exactly 'Project Scope and Objectives' (0.6 points)
    # FAILS on initial (title is 'Scope') -> PASSES on golden ('Project Scope and Objectives')
    try:
        slide2 = prs.slides[1]  # 0-indexed
        slide2_title = get_slide_title(slide2)

        if slide2_title == EXPECTED_SLIDE2_TITLE:
            print(f"PASS: Component 1 — Slide 2 title is '{slide2_title}' (0.6 pts)")
            total_score += 0.6
        else:
            print(f"FAIL: Component 1 — Expected slide 2 title '{EXPECTED_SLIDE2_TITLE}', "
                  f"found '{slide2_title}'")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 2 title is correct AND body text unchanged AND slide count is 5 (0.4 points)
    # This is a compound check anchored to the title change.
    # FAILS on initial (title check fails: 'Scope' != 'Project Scope and Objectives')
    # PASSES on golden (all three sub-conditions satisfied)
    try:
        slide2 = prs.slides[1]
        slide2_title = get_slide_title(slide2)
        title_correct = (slide2_title == EXPECTED_SLIDE2_TITLE)

        slide_count_ok = (len(prs.slides) == EXPECTED_SLIDE_COUNT)

        body_lines = get_slide_body_lines(slide2)
        body_ok = (body_lines == EXPECTED_SLIDE2_BODY_LINES)

        if title_correct and slide_count_ok and body_ok:
            print(f"PASS: Component 2 — Title correct, slide count={len(prs.slides)}, "
                  f"body intact ({len(body_lines)} lines) (0.4 pts)")
            total_score += 0.4
        else:
            if not title_correct:
                print(f"FAIL: Component 2 — Title not correct: '{slide2_title}'")
            if not slide_count_ok:
                print(f"FAIL: Component 2 — Expected {EXPECTED_SLIDE_COUNT} slides, "
                      f"found {len(prs.slides)}")
            if not body_ok:
                print(f"FAIL: Component 2 — Slide 2 body text changed or missing. "
                      f"Found {len(body_lines)} lines, expected {len(EXPECTED_SLIDE2_BODY_LINES)}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against golden file (path on VM)
file_path = f'{WORKDIR}/{TASK_ID}_initial.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
