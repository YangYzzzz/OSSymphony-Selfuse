"""
Initial Setup: Create project_proposal.pptx with 5 slides for slide title editing task.
Task ID: impress_text_039
Domain: libreoffice_impress

Creates a business project proposal presentation with 5 slides.
Slide 2 title is 'Scope' (to be changed to 'Project Scope and Objectives' in golden).
File is saved to:
  - ~/Desktop/project_proposal.pptx (the file the agent interacts with)
  - ~/impress_text_039_initial.pptx  (reference copy for reward-gen)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import shutil

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
TASK_ID = 'impress_text_039'
TASK_FILE = f'{DESKTOP}/project_proposal.pptx'
INITIAL = f'{WORKDIR}/{TASK_ID}_initial.pptx'


def add_slide_content(prs, layout_idx, title_text, body_lines):
    """Add a slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    title_shape = slide.shapes.title
    title_shape.text = title_text

    if body_lines and len(slide.placeholders) > 1:
        body_ph = slide.placeholders[1]
        tf = body_ph.text_frame
        tf.clear()
        for i, line in enumerate(body_lines):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            para.text = line
            para.level = 0

    return slide


def create_initial():
    os.makedirs(DESKTOP, exist_ok=True)

    prs = Presentation()

    # Slide 1: Title slide (layout 0)
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = 'Greenfield Infrastructure Modernization'
    slide1.placeholders[1].text = 'Q2 2025 Project Proposal\nTechnology & Operations Division'

    # Slide 2: Title + Content — title is 'Scope' (will be changed in golden)
    add_slide_content(prs, 1, 'Scope', [
        'Define the boundaries of the infrastructure modernization effort',
        'Covers data center consolidation across 3 regional offices',
        'Network redesign for improved redundancy and throughput',
        'Cloud migration of 40 legacy on-premise workloads',
        'Excludes customer-facing application redesign',
    ])

    # Slide 3: Objectives
    add_slide_content(prs, 1, 'Objectives', [
        'Reduce infrastructure operational costs by 30% within 18 months',
        'Achieve 99.95% uptime SLA for all critical systems',
        'Improve deployment velocity from bi-weekly to daily releases',
        'Eliminate single points of failure in the primary data center',
        'Standardize tooling across all engineering teams',
    ])

    # Slide 4: Timeline
    add_slide_content(prs, 1, 'Project Timeline', [
        'Phase 1 (Months 1-3): Discovery, assessment, and planning',
        'Phase 2 (Months 4-7): Network redesign and data center consolidation',
        'Phase 3 (Months 8-13): Cloud migration of legacy workloads',
        'Phase 4 (Months 14-18): Optimization, hardening, and handover',
        'Key Milestone: Production cutover targeted for Month 13',
    ])

    # Slide 5: Budget & Resources
    add_slide_content(prs, 1, 'Budget and Resources', [
        'Total estimated budget: $4.2M over 18 months',
        'Internal headcount: 8 FTE engineers, 2 project managers',
        'External vendors: Network hardware ($1.1M), Cloud credits ($850K)',
        'Contingency reserve: 15% of total project budget',
        'Steering committee review scheduled monthly',
    ])

    # Save task file at ~/Desktop/project_proposal.pptx
    prs.save(TASK_FILE)
    print(f'Task file created: {TASK_FILE}')

    # Also copy as initial reference file
    shutil.copy(TASK_FILE, INITIAL)
    print(f'Initial reference file created: {INITIAL}')

    # Verify
    prs_check = Presentation(TASK_FILE)
    print(f'Slides: {len(prs_check.slides)}')
    for i, slide in enumerate(prs_check.slides):
        title = slide.shapes.title.text if slide.shapes.title else '(no title)'
        print(f'  Slide {i+1}: {title}')


create_initial()
