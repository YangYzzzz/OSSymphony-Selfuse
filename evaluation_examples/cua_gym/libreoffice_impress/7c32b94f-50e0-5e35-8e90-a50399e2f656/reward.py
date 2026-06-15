"""
Reward Script: Reorder the first 4 slides from 1,2,3,4 to 2,4,1,3.
Task ID: osworld_impress_slide_duplication_reorder_007
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): Slide 1 is "Literature Review" AND Slide 2 is "Results"
  Component 2 (0.5): Slide 3 is "Introduction" AND Slide 4 is "Methods"
  Total: 1.0

Both components verify changes introduced by the reordering task.
Initial state: Introduction(1), Literature Review(2), Methods(3), Results(4), Discussion(5), Conclusion(6)
Golden state:  Literature Review(1), Results(2), Introduction(3), Methods(4), Discussion(5), Conclusion(6)
Slides 5 and 6 are NOT scored (they are preconditions, unchanged between initial and golden).
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_slide_duplication_reorder_007'


def get_slide_title(slide):
    """Extract the title text from a slide using placeholder idx=0 or first text shape."""
    # Try title placeholder first (idx=0)
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.idx == 0 and shape.has_text_frame:
                return shape.text.strip()
    # Fallback: first text-bearing shape
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text.strip()
            if txt:
                return txt
    return ''


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0

    Only scoring task-introduced changes (the reordering of slides 1-4).
    Slides 5 and 6 are unchanged between initial and golden, so they are NOT scored.
    """
    total_score = 0.0

    # Load presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have 6 slides
    num_slides = len(prs.slides)
    if num_slides != 6:
        print(f"CRITICAL: Expected 6 slides, found {num_slides}. Slide count mismatch.")
        print("REWARD: 0.0")
        return 0.0

    # Collect actual slide titles
    actual_titles = []
    for i, slide in enumerate(prs.slides):
        title = get_slide_title(slide)
        actual_titles.append(title)
        print(f"  Slide {i+1} title: '{title}'")

    # Component 1: Slide 1 is "Literature Review" AND Slide 2 is "Results" (0.5 points)
    # Initial state: Slide 1='Introduction', Slide 2='Literature Review'  -> FAILS on initial
    # Golden state:  Slide 1='Literature Review', Slide 2='Results'        -> PASSES on golden
    try:
        slide1_ok = actual_titles[0] == 'Literature Review'
        slide2_ok = actual_titles[1] == 'Results'
        if slide1_ok and slide2_ok:
            print(f"PASS: Component 1 — Slide 1='Literature Review', Slide 2='Results' (0.5 pts)")
            total_score += 0.5
        else:
            if not slide1_ok:
                print(f"FAIL: Component 1 — Slide 1: expected 'Literature Review', found '{actual_titles[0]}'")
            if not slide2_ok:
                print(f"FAIL: Component 1 — Slide 2: expected 'Results', found '{actual_titles[1]}'")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 3 is "Introduction" AND Slide 4 is "Methods" (0.5 points)
    # Initial state: Slide 3='Methods', Slide 4='Results'       -> FAILS on initial
    # Golden state:  Slide 3='Introduction', Slide 4='Methods'  -> PASSES on golden
    try:
        slide3_ok = actual_titles[2] == 'Introduction'
        slide4_ok = actual_titles[3] == 'Methods'
        if slide3_ok and slide4_ok:
            print(f"PASS: Component 2 — Slide 3='Introduction', Slide 4='Methods' (0.5 pts)")
            total_score += 0.5
        else:
            if not slide3_ok:
                print(f"FAIL: Component 2 — Slide 3: expected 'Introduction', found '{actual_titles[2]}'")
            if not slide4_ok:
                print(f"FAIL: Component 2 — Slide 4: expected 'Methods', found '{actual_titles[3]}'")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
