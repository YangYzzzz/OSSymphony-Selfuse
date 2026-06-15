"""
Initial Setup: 6-slide academic lecture presentation — slides in original order
Task ID: osworld_impress_slide_duplication_reorder_007
Domain: libreoffice_impress

Creates a 6-slide academic lecture presentation with slides in order:
1. Introduction
2. Literature Review
3. Methods
4. Results
5. Discussion
6. Conclusion

The agent task is to reorder the first 4 slides from [1,2,3,4] to [2,4,1,3].
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_slide_duplication_reorder_007'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_content_slide(prs, title_text, content_lines, bg_rgb=None):
    """Add a content slide with title and bullet points."""
    slide_layout = prs.slide_layouts[1]  # Title + Content layout
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title_text

    # Set background if specified
    if bg_rgb:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*bg_rgb)

    # Set content
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.level = 0

    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Introduction ---
    slide1_content = [
        "Research Topic: Impact of Climate Change on Crop Yields",
        "Investigates temperature and precipitation patterns (1990–2024)",
        "Study covers 12 agricultural regions across 4 continents",
        "Goal: quantify yield changes attributable to climate variability",
        "Funded by the Global Agricultural Research Consortium (GARC)",
    ]
    s1 = add_content_slide(prs, "Introduction", slide1_content)
    s1.notes_slide.notes_text_frame.text = "Slide 1: Introduction — overview of the study"

    # --- Slide 2: Literature Review ---
    slide2_content = [
        "Prior studies (Smith et al., 2018; Patel & Wu, 2020) show 2–8% yield decline per °C",
        "Hansen (2021) projected a 15% reduction in wheat yields by 2050",
        "Nguyen et al. (2022): irrigation mitigates up to 40% of heat stress losses",
        "Gap identified: limited data from sub-Saharan Africa and Southeast Asia",
        "Our study addresses this gap with new field data (2019–2024)",
    ]
    s2 = add_content_slide(prs, "Literature Review", slide2_content)
    s2.notes_slide.notes_text_frame.text = "Slide 2: Literature Review — prior work"

    # --- Slide 3: Methods ---
    slide3_content = [
        "Data collection: 4,200 field stations; daily temp, rainfall, soil moisture",
        "Satellite imagery (Sentinel-2, MODIS) for crop-area classification",
        "Statistical model: panel regression with fixed effects by region",
        "Validation: cross-validation on held-out 20% of station data (R² = 0.87)",
        "Software: Python 3.11, scikit-learn, xarray; code available on GitHub",
    ]
    s3 = add_content_slide(prs, "Methods", slide3_content)
    s3.notes_slide.notes_text_frame.text = "Slide 3: Methods — data collection and analysis pipeline"

    # --- Slide 4: Results ---
    slide4_content = [
        "Maize yields declined by 6.3% per °C increase (p < 0.001)",
        "Rice yields declined by 4.1% per °C; wheat by 7.8% per °C",
        "Sub-Saharan Africa showed largest impact: −11.2% per °C for sorghum",
        "Irrigation-equipped farms showed 35% smaller yield losses on average",
        "Total projected loss by 2040: 220 million metric tons globally (base scenario)",
    ]
    s4 = add_content_slide(prs, "Results", slide4_content)
    s4.notes_slide.notes_text_frame.text = "Slide 4: Results — key findings"

    # --- Slide 5: Discussion ---
    slide5_content = [
        "Findings align with and extend Hansen (2021) projections",
        "Heat-tolerant cultivar adoption could offset up to 50% of projected losses",
        "Policy implication: prioritize irrigation infrastructure in high-risk regions",
        "Limitations: model does not capture CO₂ fertilization effects fully",
        "Future work: incorporate socioeconomic adaptation pathways",
    ]
    s5 = add_content_slide(prs, "Discussion", slide5_content)
    s5.notes_slide.notes_text_frame.text = "Slide 5: Discussion — interpretation and implications"

    # --- Slide 6: Conclusion ---
    slide6_content = [
        "Climate change poses a measurable and growing threat to global food security",
        "Yield losses of 4–12% per °C are consistent across crop types and regions",
        "Targeted investment in irrigation and heat-tolerant varieties is essential",
        "New dataset released publicly: https://garc-data.org/climate-crop-2025",
        "Thank you — Questions welcome",
    ]
    s6 = add_content_slide(prs, "Conclusion", slide6_content)
    s6.notes_slide.notes_text_frame.text = "Slide 6: Conclusion — summary and call to action"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print('Slide order (initial):')
    prs2 = Presentation(OUTPUT)
    for i, sl in enumerate(prs2.slides):
        title = sl.shapes.title.text if sl.shapes.title else '(no title)'
        print(f'  Slide {i+1}: {title}')

    # GUI-ready startup: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
