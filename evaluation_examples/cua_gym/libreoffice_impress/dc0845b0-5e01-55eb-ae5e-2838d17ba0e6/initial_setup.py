"""
Initial Setup: Blank presentation for DevConf talk creation task
Task ID: impress_wf_044
Domain: libreoffice_impress

Creates a blank presentation and opens it in LibreOffice Impress.
The agent's task is to build the full 8-slide conference presentation from scratch.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_044'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    # Create a blank presentation (the context says Impress is open with a blank presentation)
    prs = Presentation()
    # Add one blank slide so it's not completely empty
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only layout, blank-ish

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
