"""
Reward Script: Conference presentation with speaker slides (DevConf 2024)
Task ID: impress_wf_044
Domain: libreoffice_impress
Scoring:
  C1 (0.15): File on Desktop with 8 slides
  C2 (0.15): All 8 slides have #263238 background
  C3 (0.10): Slide 1 has title 'DevConf 2024 - Building Scalable APIs'
  C4 (0.10): Slide 2 has a circle/oval shape (headshot placeholder)
  C5 (0.20): Slides 4-6 have rounded-rect code blocks with monospace font
  C6 (0.10): Slides 4-6 have Appear entrance animation
  C7 (0.10): Slide 7 has black rectangle + 'DEMO' text
  C8 (0.10): Slide 8 has '?' text and social handle
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_044'
FILE_PATH = os.path.join(WORKDIR, 'Desktop', 'DevConf_Talk.pptx')


def get_slide_bg_rgb(slide):
    """Return background RGB as uppercase hex string or None."""
    fill = slide.background.fill
    if fill.type == 1:  # SOLID
        return str(fill.fore_color.rgb).upper()
    elif fill.type == 5:  # BACKGROUND (inherited)
        try:
            master_fill = slide.slide_layout.slide_master.background.fill
            if master_fill.type == 1:
                return str(master_fill.fore_color.rgb).upper()
        except Exception:
            pass
    return None


def get_all_text(slide):
    """Collect all text from a slide's shapes (including groups)."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return texts


def has_appear_animation(pptx_path, slide_num):
    """Check if slide has Appear entrance animation (presetID=1, presetClass=entr)."""
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            fname = f'ppt/slides/slide{slide_num}.xml'
            with zf.open(fname) as f:
                root = ET.parse(f).getroot()
                # Find all cTn elements with presetID=1 and presetClass=entr
                for elem in root.iter():
                    tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
                    if tag == 'cTn':
                        pid = elem.get('presetID')
                        pcl = elem.get('presetClass')
                        if pid == '1' and pcl == 'entr':
                            return True
        return False
    except Exception:
        return False


def verify_task(file_path):
    """Verify task completion with progressive scoring. Returns float 0.0-1.0."""
    total_score = 0.0

    # Load presentation
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)
    num_slides = len(slides)

    # Component 1: File on Desktop with 8 slides (0.15 points)
    try:
        if num_slides == 8:
            print(f"PASS: Component 1 — File has {num_slides} slides (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — Expected 8 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: All 8 slides have #263238 background (0.15 points)
    try:
        if num_slides >= 8:
            bg_pass = 0
            for i in range(8):
                bg = get_slide_bg_rgb(slides[i])
                if bg == '263238':
                    bg_pass += 1
                else:
                    print(f"  Slide {i+1} background: {bg} (expected 263238)")
            if bg_pass == 8:
                print(f"PASS: Component 2 — All 8 slides have #263238 background (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 2 — {bg_pass}/8 slides have correct background")
        else:
            print(f"FAIL: Component 2 — Not enough slides to check backgrounds")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 1 has title 'DevConf 2024 - Building Scalable APIs' (0.10 points)
    try:
        if num_slides >= 1:
            slide1_texts = get_all_text(slides[0])
            full_text = ' '.join(slide1_texts).lower()
            if 'devconf 2024' in full_text and 'building scalable apis' in full_text:
                print(f"PASS: Component 3 — Slide 1 has conference title (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 3 — Slide 1 texts: {slide1_texts[:3]}")
        else:
            print(f"FAIL: Component 3 — No slides")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 2 has a circle/oval shape for headshot placeholder (0.10 points)
    try:
        if num_slides >= 2:
            has_oval = False
            for shape in slides[1].shapes:
                try:
                    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
                    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        ast = shape.auto_shape_type
                        # OVAL = 9, could also be ELLIPSE
                        if ast is not None and ast in (9,):  # MSO_AUTO_SHAPE_TYPE.OVAL
                            has_oval = True
                            break
                except Exception:
                    pass
            if has_oval:
                print(f"PASS: Component 4 — Slide 2 has oval/circle shape (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 4 — No oval/circle shape found on slide 2")
        else:
            print(f"FAIL: Component 4 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slides 4-6 have rounded-rect code blocks with monospace font (0.20 points)
    try:
        if num_slides >= 6:
            code_slides_pass = 0
            for idx in [3, 4, 5]:  # 0-indexed slides 4, 5, 6
                slide = slides[idx]
                has_code_block = False
                for shape in slide.shapes:
                    try:
                        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                            ast = shape.auto_shape_type
                            # ROUNDED_RECTANGLE = 5
                            if ast is not None and ast == 5:
                                # Check for monospace font in text
                                if shape.has_text_frame:
                                    for para in shape.text_frame.paragraphs:
                                        for run in para.runs:
                                            if run.text.strip():
                                                fn = (run.font.name or '').lower()
                                                if fn in ('courier new', 'courier', 'consolas', 'monospace',
                                                          'dejavu sans mono', 'liberation mono', 'source code pro'):
                                                    has_code_block = True
                                                    break
                                        if has_code_block:
                                            break
                    except Exception:
                        pass
                    if has_code_block:
                        break
                if has_code_block:
                    code_slides_pass += 1
                else:
                    print(f"  Slide {idx+1}: no rounded-rect with monospace font found")

            if code_slides_pass == 3:
                print(f"PASS: Component 5 — All 3 slides (4-6) have code blocks (0.20 pts)")
                total_score += 0.20
            elif code_slides_pass >= 1:
                partial = round(0.20 * code_slides_pass / 3, 2)
                print(f"PARTIAL: Component 5 — {code_slides_pass}/3 slides have code blocks ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 5 — No code blocks found on slides 4-6")
        else:
            print(f"FAIL: Component 5 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Slides 4-6 have Appear entrance animation (0.10 points)
    try:
        if num_slides >= 6:
            anim_pass = 0
            for slide_num in [4, 5, 6]:  # 1-indexed
                if has_appear_animation(file_path, slide_num):
                    anim_pass += 1
                else:
                    print(f"  Slide {slide_num}: no Appear animation found")
            if anim_pass == 3:
                print(f"PASS: Component 6 — All 3 slides (4-6) have Appear animation (0.10 pts)")
                total_score += 0.10
            elif anim_pass >= 1:
                partial = round(0.10 * anim_pass / 3, 2)
                print(f"PARTIAL: Component 6 — {anim_pass}/3 slides have animation ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 6 — No Appear animations on slides 4-6")
        else:
            print(f"FAIL: Component 6 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Slide 7 has black rectangle + 'DEMO' text (0.10 points)
    try:
        if num_slides >= 7:
            slide7 = slides[6]
            has_black_rect = False
            has_demo_text = False

            for shape in slide7.shapes:
                # Check for black filled rectangle
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        sf = shape.fill
                        if sf.type is not None:
                            try:
                                fill_rgb = str(sf.fore_color.rgb).upper()
                                if fill_rgb == '000000':
                                    has_black_rect = True
                            except Exception:
                                pass
                except Exception:
                    pass

                # Check for DEMO text
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if 'DEMO' in para.text.upper():
                            has_demo_text = True

            if has_black_rect and has_demo_text:
                print(f"PASS: Component 7 — Slide 7 has black rectangle + DEMO text (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 7 — black_rect={has_black_rect}, demo_text={has_demo_text}")
        else:
            print(f"FAIL: Component 7 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    # Component 8: Slide 8 has '?' text and social handle (0.10 points)
    try:
        if num_slides >= 8:
            slide8 = slides[7]
            slide8_texts = get_all_text(slide8)
            full_text = ' '.join(slide8_texts)
            has_question = '?' in full_text
            has_social = '@' in full_text  # social handle like @elena_api

            if has_question and has_social:
                print(f"PASS: Component 8 — Slide 8 has '?' and social handle (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 8 — question_mark={has_question}, social={has_social}")
        else:
            print(f"FAIL: Component 8 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 8 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(FILE_PATH)
