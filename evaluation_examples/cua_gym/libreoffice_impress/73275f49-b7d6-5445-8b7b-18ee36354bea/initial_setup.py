"""
Initial Setup: Center-align body text and set paragraph spacing on slide 4
Task ID: impress_tct_073
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_073'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_placeholder(placeholder, texts, font_size=Pt(18), bold=False, alignment=PP_ALIGN.LEFT):
    """Add multiple paragraphs of text to a placeholder."""
    tf = placeholder.text_frame
    tf.word_wrap = True
    for i, text in enumerate(texts):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.alignment = alignment
        for run in p.runs:
            run.font.size = font_size
            run.font.bold = bold


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Annual Tech Conference 2025"
    slide1.placeholders[1].text = "Innovation Through Collaboration"

    # --- Slide 2: Keynote Speakers ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Keynote Speakers"
    body2 = slide2.placeholders[1]
    add_text_to_placeholder(body2, [
        "Dr. Elena Rodriguez - AI Ethics in Healthcare",
        "James Park - The Future of Quantum Computing",
        "Priya Sharma - Sustainable Tech Infrastructure",
        "Michael Chen - Cybersecurity in the Age of AI",
    ])

    # --- Slide 3: Workshop Schedule ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Workshop Schedule"
    body3 = slide3.placeholders[1]
    add_text_to_placeholder(body3, [
        "Morning Session (9:00 - 12:00): Hands-on Machine Learning",
        "Lunch Break (12:00 - 13:00): Networking Lounge Open",
        "Afternoon Session A (13:00 - 15:00): Cloud Architecture Patterns",
        "Afternoon Session B (15:30 - 17:00): DevOps Best Practices",
    ])

    # --- Slide 4: Venue & Logistics (TARGET SLIDE) ---
    # LEFT-aligned, default spacing (NO center, NO custom spacing)
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Venue & Logistics"
    body4 = slide4.placeholders[1]
    tf4 = body4.text_frame
    tf4.word_wrap = True

    paragraphs_text = [
        "The conference will be held at the Grand Pacific Convention Center, located in downtown San Francisco at 450 Mission Street.",
        "Complimentary shuttle service runs every 15 minutes from Union Square and the Moscone Center area between 7:30 AM and 8:00 PM daily.",
        "On-site parking is available in the underground garage for $25 per day, with early-bird rates of $18 for arrivals before 8:00 AM.",
        "Lunch and refreshments are provided for all registered attendees. Please indicate any dietary restrictions during online check-in.",
    ]

    for i, text in enumerate(paragraphs_text):
        if i == 0:
            p = tf4.paragraphs[0]
        else:
            p = tf4.add_paragraph()
        p.text = text
        p.alignment = PP_ALIGN.LEFT
        # Default spacing: no space_before / space_after set
        for run in p.runs:
            run.font.size = Pt(16)

    # --- Slide 5: Sponsorship Tiers ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Sponsorship Tiers"
    body5 = slide5.placeholders[1]
    add_text_to_placeholder(body5, [
        "Platinum ($50,000): Main stage branding, 10 VIP passes, booth",
        "Gold ($25,000): Session sponsorship, 5 VIP passes, booth",
        "Silver ($10,000): Logo on materials, 3 standard passes",
        "Bronze ($5,000): Website listing, 2 standard passes",
    ])

    # --- Slide 6: Contact & Registration ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Contact & Registration"
    body6 = slide6.placeholders[1]
    add_text_to_placeholder(body6, [
        "Register at: www.techconf2025.io/register",
        "Early bird pricing ends March 15, 2025",
        "Questions: events@techconf2025.io or (415) 555-0192",
        "Follow us: @TechConf2025 on all platforms",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
