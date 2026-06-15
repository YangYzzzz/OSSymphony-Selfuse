"""
Reward Script: Center-align body text on slide 4 and set paragraph spacing
Task ID: impress_tct_073
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): All 4 body paragraphs on slide 4 are center-aligned
  Component 2 (0.25): All 4 body paragraphs have space_before = 12pt (152400 EMU)
  Component 3 (0.25): All 4 body paragraphs have space_after = 6pt (76200 EMU)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_073'
TARGET_FILE = f'{WORKDIR}/{TASK_ID}.pptx'

# Expected values
EXPECTED_ALIGN = 2  # PP_ALIGN.CENTER
EXPECTED_SPACE_BEFORE = 152400  # 12pt in EMU
EXPECTED_SPACE_AFTER = 76200    # 6pt in EMU
TARGET_SLIDE_IDX = 3  # Slide 4 (0-indexed)
CONTENT_SHAPE_NAME = "Content Placeholder 2"
EXPECTED_PARA_COUNT = 4


def persist_app_state():
    """Save any unsaved LibreOffice edits before verification."""
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.text import PP_ALIGN
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Presentation has at least 4 slides
    if len(prs.slides) < 4:
        print(f"PRECONDITION FAIL: Expected at least 4 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[TARGET_SLIDE_IDX]

    # Find the content text box (not the title)
    content_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == CONTENT_SHAPE_NAME:
            content_shape = shape
            break

    if content_shape is None:
        # Fallback: find a non-title text shape with multiple paragraphs
        for shape in slide.shapes:
            if shape.has_text_frame:
                paras = [p for p in shape.text_frame.paragraphs if p.text.strip()]
                if len(paras) >= EXPECTED_PARA_COUNT:
                    content_shape = shape
                    break

    if content_shape is None:
        print("PRECONDITION FAIL: No content text box found on slide 4 with 4+ paragraphs")
        print("REWARD: 0.0")
        return 0.0

    paragraphs = [p for p in content_shape.text_frame.paragraphs if p.text.strip()]
    print(f"INFO: Found {len(paragraphs)} non-empty paragraphs in '{content_shape.name}'")

    # Component 1: All body paragraphs are center-aligned (0.5 points)
    try:
        centered_count = 0
        for i, para in enumerate(paragraphs):
            align_val = para.alignment
            # PP_ALIGN.CENTER == 2
            if align_val is not None and int(align_val) == EXPECTED_ALIGN:
                centered_count += 1
                print(f"  Para {i}: align=CENTER -- OK")
            else:
                print(f"  Para {i}: align={align_val} -- expected CENTER (2)")

        if centered_count == len(paragraphs) and len(paragraphs) >= EXPECTED_PARA_COUNT:
            print(f"PASS: Component 1 -- All {centered_count} paragraphs center-aligned (0.5 pts)")
            total_score += 0.5
        elif centered_count > 0:
            partial = 0.5 * (centered_count / len(paragraphs))
            print(f"PARTIAL: Component 1 -- {centered_count}/{len(paragraphs)} centered ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- No paragraphs are center-aligned")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Space before each paragraph is 12pt / 152400 EMU (0.25 points)
    try:
        correct_before = 0
        for i, para in enumerate(paragraphs):
            spc = para.space_before
            if spc is not None and int(spc) == EXPECTED_SPACE_BEFORE:
                correct_before += 1
                print(f"  Para {i}: space_before={spc} (12pt) -- OK")
            else:
                print(f"  Para {i}: space_before={spc} -- expected {EXPECTED_SPACE_BEFORE} (12pt)")

        if correct_before == len(paragraphs) and len(paragraphs) >= EXPECTED_PARA_COUNT:
            print(f"PASS: Component 2 -- All {correct_before} paragraphs have 12pt space before (0.25 pts)")
            total_score += 0.25
        elif correct_before > 0:
            partial = 0.25 * (correct_before / len(paragraphs))
            print(f"PARTIAL: Component 2 -- {correct_before}/{len(paragraphs)} correct ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 -- No paragraphs have correct space before")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Space after each paragraph is 6pt / 76200 EMU (0.25 points)
    try:
        correct_after = 0
        for i, para in enumerate(paragraphs):
            spc = para.space_after
            if spc is not None and int(spc) == EXPECTED_SPACE_AFTER:
                correct_after += 1
                print(f"  Para {i}: space_after={spc} (6pt) -- OK")
            else:
                print(f"  Para {i}: space_after={spc} -- expected {EXPECTED_SPACE_AFTER} (6pt)")

        if correct_after == len(paragraphs) and len(paragraphs) >= EXPECTED_PARA_COUNT:
            print(f"PASS: Component 3 -- All {correct_after} paragraphs have 6pt space after (0.25 pts)")
            total_score += 0.25
        elif correct_after > 0:
            partial = 0.25 * (correct_after / len(paragraphs))
            print(f"PARTIAL: Component 3 -- {correct_after}/{len(paragraphs)} correct ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 -- No paragraphs have correct space after")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

if not os.path.exists(TARGET_FILE):
    print(f"File not found: {TARGET_FILE}")
    print("REWARD: 0.0")
else:
    verify_task(TARGET_FILE)
