"""
Reward Script: Insert a complex table on slide 5 with merged cells for Weekly Study Plan
Task ID: impress_stu_067
Domain: libreoffice_impress
Scoring:
  1. Table exists on slide 5 with correct dimensions (0.15)
  2. Header row "Week 5 Study Schedule" spanning full width (0.15)
  3. Sub-header row content: Time, Monday, Wednesday, Friday (0.15)
  4. Data rows content correct (0.15)
  5. Reading cells merged (gridSpan=2) (0.15)
  6. Header styling: blue bg #3498DB, white bold text (0.15)
  7. Alternating row colors: white/gray (0.10)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_067'

NS = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
      'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
      'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}


def persist_app_state(domain):
    """Best-effort save via Ctrl+S for LibreOffice."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def get_cell_text(tc, ns):
    """Extract all text from a table cell XML element."""
    return ''.join(t.text or '' for t in tc.findall('.//a:t', ns))


def get_cell_fill_color(tc, ns):
    """Extract solid fill color from a table cell XML element. Returns hex string or None."""
    # Look for solidFill in tcPr
    tcPr = tc.find('a:tcPr', ns)
    if tcPr is None:
        return None
    solidFill = tcPr.find('a:solidFill', ns)
    if solidFill is None:
        return None
    srgbClr = solidFill.find('a:srgbClr', ns)
    if srgbClr is not None:
        return srgbClr.get('val', '').upper()
    return None


def get_run_font_color(tc, ns):
    """Extract font color from a run in a cell. Returns hex string or None."""
    for rPr in tc.findall('.//a:rPr', ns):
        solidFill = rPr.find('a:solidFill', ns)
        if solidFill is not None:
            srgbClr = solidFill.find('a:srgbClr', ns)
            if srgbClr is not None:
                return srgbClr.get('val', '').upper()
    return None


def is_run_bold(tc, ns):
    """Check if any run in the cell has bold text."""
    for rPr in tc.findall('.//a:rPr', ns):
        b = rPr.get('b')
        if b == '1' or b == 'true':
            return True
    return False


def verify_task(file_path):
    """Verify task completion with progressive scoring. Returns float 0.0..1.0."""
    total_score = 0.0

    # We need both python-pptx (for basic checks) and XML (for merge/style checks)
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # slide 5, 0-indexed

    # Component 1: Table exists on slide 5 with correct dimensions (0.15 points)
    try:
        table_shape = None
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_shape = shape
                break

        if table_shape is not None:
            table = table_shape.table
            rows = len(table.rows)
            cols = len(table.columns)
            if rows == 5 and cols == 4:
                print(f"PASS: Component 1 -- Table on slide 5 with 5x4 dimensions (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 1 -- Table dimensions {rows}x{cols}, expected 5x4")
        else:
            print(f"FAIL: Component 1 -- No table found on slide 5")
            print(f"REWARD: 0.0")
            return 0.0
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Parse XML for detailed checks
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide5.xml') as f:
                slide_xml = ET.parse(f).getroot()
        tbl = slide_xml.find('.//a:tbl', NS)
        if tbl is None:
            print("FAIL: No table found in slide5.xml")
            print(f"REWARD: {total_score}")
            return total_score
        trs = tbl.findall('a:tr', NS)
    except Exception as e:
        print(f"ERROR: Cannot parse slide5.xml: {e}")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Header row "Week 5 Study Schedule" spanning full width (0.15 points)
    try:
        row0_cells = trs[0].findall('a:tc', NS)
        header_text = get_cell_text(row0_cells[0], NS).strip()
        gridSpan = row0_cells[0].get('gridSpan')

        header_text_ok = 'week 5 study schedule' in header_text.lower()
        header_span_ok = gridSpan is not None and int(gridSpan) == 4

        if header_text_ok and header_span_ok:
            print(f"PASS: Component 2 -- Header '{header_text}' spans 4 columns (0.15 pts)")
            total_score += 0.15
        elif header_text_ok:
            print(f"PARTIAL: Component 2 -- Header text correct but gridSpan={gridSpan} (0.07 pts)")
            total_score += 0.07
        elif header_span_ok:
            print(f"PARTIAL: Component 2 -- Header spans 4 cols but text='{header_text}' (0.07 pts)")
            total_score += 0.07
        else:
            print(f"FAIL: Component 2 -- Header text='{header_text}', gridSpan={gridSpan}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Sub-header row content (0.15 points)
    try:
        expected_subheaders = ['Time', 'Monday', 'Wednesday', 'Friday']
        row1_cells = trs[1].findall('a:tc', NS)
        actual_subheaders = [get_cell_text(c, NS).strip() for c in row1_cells]

        matches = sum(1 for exp, act in zip(expected_subheaders, actual_subheaders)
                      if exp.lower() == act.lower())

        if matches == 4:
            print(f"PASS: Component 3 -- Sub-headers correct: {actual_subheaders} (0.15 pts)")
            total_score += 0.15
        elif matches >= 2:
            pts = round(0.15 * matches / 4, 2)
            print(f"PARTIAL: Component 3 -- {matches}/4 sub-headers correct ({pts} pts)")
            total_score += pts
        else:
            print(f"FAIL: Component 3 -- Sub-headers: {actual_subheaders}, expected {expected_subheaders}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Data rows content correct (0.15 points)
    try:
        expected_data = {
            2: ['9-11 AM', 'Calculus review', 'Physics lab prep', 'Calculus review'],
            3: ['11-1 PM', 'Reading', '', 'Essay writing'],  # col 2 may be empty due to merge
            4: ['2-4 PM', 'Group study', 'Office hours', 'Practice problems'],
        }

        data_matches = 0
        total_checks = 0

        for row_idx in [2, 3, 4]:
            row_cells = trs[row_idx].findall('a:tc', NS)
            expected = expected_data[row_idx]
            for ci, exp in enumerate(expected):
                if ci < len(row_cells):
                    actual = get_cell_text(row_cells[ci], NS).strip()
                    # For merged reading cell, accept either text in merged or empty in continuation
                    if row_idx == 3 and ci == 2:
                        # This cell may be merged continuation (hMerge) so empty is OK
                        # or may contain "Reading" if merged differently
                        if actual == '' or actual.lower() == 'reading':
                            data_matches += 1
                    elif exp == '':
                        if actual == '':
                            data_matches += 1
                    else:
                        if exp.lower() == actual.lower():
                            data_matches += 1
                        else:
                            print(f"  Data mismatch row {row_idx}, col {ci}: expected '{exp}', got '{actual}'")
                total_checks += 1

        if total_checks > 0:
            ratio = data_matches / total_checks
            pts = round(0.15 * ratio, 2)
            if ratio >= 0.9:
                print(f"PASS: Component 4 -- Data rows correct ({data_matches}/{total_checks}) (0.15 pts)")
                total_score += 0.15
            elif ratio > 0:
                print(f"PARTIAL: Component 4 -- {data_matches}/{total_checks} cells correct ({pts} pts)")
                total_score += pts
            else:
                print(f"FAIL: Component 4 -- No data cells match")
        else:
            print(f"FAIL: Component 4 -- No data rows found")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Reading cells merged (gridSpan=2) (0.15 points)
    try:
        row3_cells = trs[3].findall('a:tc', NS)
        # Cell at col 1 (Monday) should have gridSpan=2 for merged "Reading"
        reading_cell = row3_cells[1] if len(row3_cells) > 1 else None
        reading_text = get_cell_text(reading_cell, NS).strip() if reading_cell is not None else ''
        reading_span = reading_cell.get('gridSpan') if reading_cell is not None else None

        merge_ok = reading_span is not None and int(reading_span) >= 2
        text_ok = 'reading' in reading_text.lower()

        if merge_ok and text_ok:
            print(f"PASS: Component 5 -- Reading cell merged (gridSpan={reading_span}) with text '{reading_text}' (0.15 pts)")
            total_score += 0.15
        elif merge_ok:
            print(f"PARTIAL: Component 5 -- Merge OK but text='{reading_text}' (0.07 pts)")
            total_score += 0.07
        elif text_ok:
            print(f"PARTIAL: Component 5 -- Text OK but no merge (gridSpan={reading_span}) (0.07 pts)")
            total_score += 0.07
        else:
            print(f"FAIL: Component 5 -- Reading merge: text='{reading_text}', gridSpan={reading_span}")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # Component 6: Header styling - blue bg #3498DB, white bold text (0.15 points)
    try:
        style_score = 0.0
        # Check row 0 (main header) and row 1 (sub-headers)
        header_rows = [trs[0], trs[1]]
        blue_cells = 0
        white_bold_cells = 0
        total_header_cells = 0

        for tr in header_rows:
            for tc in tr.findall('a:tc', NS):
                hMerge = tc.get('hMerge')
                if hMerge == '1':
                    continue  # skip merged continuation cells
                total_header_cells += 1
                fill_color = get_cell_fill_color(tc, NS)
                if fill_color and fill_color == '3498DB':
                    blue_cells += 1

                text = get_cell_text(tc, NS).strip()
                if text:  # only check font on cells with text
                    font_color = get_run_font_color(tc, NS)
                    bold = is_run_bold(tc, NS)
                    if font_color == 'FFFFFF' and bold:
                        white_bold_cells += 1

        cells_with_text = 5  # 1 header + 4 sub-headers
        if total_header_cells > 0:
            bg_ratio = blue_cells / total_header_cells
            text_ratio = white_bold_cells / cells_with_text if cells_with_text > 0 else 0

            if bg_ratio >= 0.8 and text_ratio >= 0.8:
                print(f"PASS: Component 6 -- Header styling correct: {blue_cells} blue cells, {white_bold_cells} white-bold cells (0.15 pts)")
                total_score += 0.15
            elif bg_ratio >= 0.5 or text_ratio >= 0.5:
                pts = round(0.15 * (bg_ratio + text_ratio) / 2, 2)
                print(f"PARTIAL: Component 6 -- Blue: {blue_cells}/{total_header_cells}, White-bold: {white_bold_cells}/{cells_with_text} ({pts} pts)")
                total_score += pts
            else:
                print(f"FAIL: Component 6 -- Blue: {blue_cells}/{total_header_cells}, White-bold: {white_bold_cells}/{cells_with_text}")
        else:
            print(f"FAIL: Component 6 -- No header cells found")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    # Component 7: Alternating row colors (0.10 points)
    try:
        # Data rows: row 2 (white), row 3 (gray), row 4 (white)
        # We check that rows alternate between two different colors
        data_rows = [trs[2], trs[3], trs[4]]
        row_colors = []
        for tr in data_rows:
            tcs = tr.findall('a:tc', NS)
            # Get color of first non-merged cell
            color = None
            for tc in tcs:
                if tc.get('hMerge') != '1':
                    color = get_cell_fill_color(tc, NS)
                    break
            row_colors.append(color)

        # Check alternation: row2 != row3 and row2 == row4
        if len(row_colors) == 3 and all(c is not None for c in row_colors):
            alternating = (row_colors[0] != row_colors[1]) and (row_colors[0] == row_colors[2])
            if alternating:
                print(f"PASS: Component 7 -- Alternating row colors: {row_colors} (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 7 -- Row colors not alternating: {row_colors}")
        else:
            print(f"FAIL: Component 7 -- Could not determine row colors: {row_colors}")
    except Exception as e:
        print(f"ERROR: Component 7 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
