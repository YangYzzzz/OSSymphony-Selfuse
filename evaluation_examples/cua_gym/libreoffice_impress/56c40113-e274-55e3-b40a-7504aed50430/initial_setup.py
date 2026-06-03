"""
Initial Setup: Insert a complex table on slide 5 with merged cells
Task ID: impress_stu_067
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_067'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_title_text(slide, text):
    """Set the title of a slide."""
    if slide.shapes.title:
        slide.shapes.title.text = text


def add_body_text(slide, texts):
    """Add body text to a content placeholder."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:  # body placeholder
            tf = ph.text_frame
            tf.clear()
            for i, txt in enumerate(texts):
                if i == 0:
                    tf.paragraphs[0].text = txt
                else:
                    p = tf.add_paragraph()
                    p.text = txt
            return


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Weekly Study Plan"
    slide1.placeholders[1].text = "Spring 2025 Semester - Academic Planning"

    # --- Slide 2: Introduction ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    set_title_text(slide2, "Introduction")
    add_body_text(slide2, [
        "This presentation outlines the weekly study plan for the spring semester.",
        "The goal is to allocate time effectively across all courses.",
        "Key focus areas: Calculus, Physics, English Literature, and Computer Science.",
        "Study sessions are organized by priority and assignment deadlines.",
    ])

    # --- Slide 3: Course Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    set_title_text(slide3, "Course Overview")
    add_body_text(slide3, [
        "MATH 201 - Calculus II: Integration techniques, series, applications",
        "PHYS 101 - General Physics: Mechanics, thermodynamics, waves",
        "ENG 150 - English Literature: Essay writing, critical analysis",
        "CS 110 - Intro to Programming: Python fundamentals, algorithms",
    ])

    # --- Slide 4: Study Resources ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    set_title_text(slide4, "Study Resources")
    add_body_text(slide4, [
        "University Library - Open 7 AM to 11 PM weekdays",
        "Online Tutoring Platform - Available 24/7 via student portal",
        "Professor Office Hours - Check department website for schedules",
        "Study Group Meetups - Science Building Room 204, Wednesdays 5 PM",
        "Practice Problem Database - Access through course LMS",
    ])

    # --- Slide 5: Study Schedule (EMPTY - task target) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add only a title text box
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Study Schedule"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    # NO table here - that's the task

    # --- Slide 6: Summary ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    set_title_text(slide6, "Summary & Next Steps")
    add_body_text(slide6, [
        "Follow the weekly schedule consistently for best results.",
        "Adjust study times based on upcoming exams and deadlines.",
        "Review progress at the end of each week.",
        "Reach out to professors or tutors when struggling with material.",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
