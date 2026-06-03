"""
FINAL REWARD SCRIPT - SUCCESS
Task: In LibreOffice Impress, slide 255 looks a bit mismatched: the bullets are still default black while the title text is that specific #0047AB blue. What’s the quickest way to grab that exact #0047AB with the eyedropper and recolor every bullet on that slide so they match perfectly?
Generated: 2025-09-10 19:32:21
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.dml.color import RGBColor


def verify_slide_bullet_color(file_path: str,
                              slide_number: int = 255,
                              expected_hex: str = "0047AB") -> float:
    """Verify that all bullet runs on the given slide use the same colour as
    the title (expected_hex) and return a progressive score between 0-1.

    Scoring rubric (max 1.0):
    • 0.4  – Title colour matches expected hex value
    • 0.2  – At least one bullet matches the title colour
    • 0.4  – ALL bullet runs match the title colour
    (Prerequisites such as file existence or slide count give NO points.)
    """

    print(">>> Loading presentation:", file_path)
    if not os.path.exists(file_path):
        print(">>> File not found – task failed.")
        print("REWARD: 0.0")
        return 0.0

    # --------------- Load presentation --------------- #
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f">>> Error loading presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    if len(prs.slides) < slide_number:
        print(f">>> Slide {slide_number} does not exist (found {len(prs.slides)} slides).")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[slide_number - 1]
    print(f">>> Analyzing slide {slide_number} (index {slide_number - 1})")

    # --------------- Identify title colour --------------- #
    reference_color = None
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        is_title = False
        if shape.is_placeholder:
            phf = shape.placeholder_format
            is_title = phf.type == PP_PLACEHOLDER.TITLE
        # Fallback on name containing "title"
        if is_title or (getattr(shape, "name", "").lower().startswith("title")):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.color.rgb is not None:
                        reference_color = run.font.color.rgb
                        break
                if reference_color is not None:
                    break
        if reference_color is not None:
            break

    if reference_color is None:
        print(">>> Could not determine title colour – no points awarded for colour match.")

    score = 0.0  # progressive scoring starts at 0

    # Compare title colour to expected hex
    try:
        expected_rgb = RGBColor.from_string(expected_hex)
    except Exception:
        expected_rgb = None

    if reference_color is not None and expected_rgb is not None and reference_color == expected_rgb:
        score += 0.4
        print(f">>> Title colour matches expected #{expected_hex} (+0.4)")
    else:
        print(f">>> Title colour does NOT match expected #{expected_hex}")

    # --------------- Evaluate bullet colours --------------- #
    total_bullet_runs = 0
    matching_bullet_runs = 0

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # Skip the title placeholder
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
            continue

        for paragraph in shape.text_frame.paragraphs:
            if not paragraph.text.strip():
                continue  # ignore empty paragraphs
            for run in paragraph.runs:
                if not run.text.strip():
                    continue
                total_bullet_runs += 1
                if reference_color is not None and run.font.color.rgb == reference_color:
                    matching_bullet_runs += 1

    print(f">>> Total bullet text runs: {total_bullet_runs}")
    print(f">>> Bullet runs matching title colour: {matching_bullet_runs}")

    if total_bullet_runs > 0:
        if matching_bullet_runs > 0:
            score += 0.2
            print(">>> At least one bullet matches title colour (+0.2)")
        if matching_bullet_runs == total_bullet_runs:
            score += 0.4
            print(">>> ALL bullets match title colour (+0.4)")
        else:
            print(">>> Not all bullets use the correct colour – partial points awarded if applicable.")
    else:
        print(">>> No bullet text runs detected on this slide (no bullet-related points).")

    final_score = min(score, 1.0)
    print(">>> Final score:", final_score)
    print(f"REWARD: {final_score}")
    return final_score


if __name__ == "__main__":
    # Path to the golden answer file within the VM
    FILE_PATH = "/home/user/in_libreoffice_impress_slide_255_looks_a_bit_mismatched_the_bullets_are_still_default_black_while_th_golden.pptx"
    verify_slide_bullet_color(FILE_PATH)

