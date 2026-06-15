"""
Initial Setup: Biology lecture presentation with 10 slides, slide 10 titled 'Key Takeaways' but empty.
Task ID: impress_teach_067
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_067'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide layouts: 0=Title, 1=Title+Content, 5=Blank, 6=Title Only
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]
    title_only_layout = prs.slide_layouts[5]  # blank for title-only effect

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Introduction to Biology"
    slide1.placeholders[1].text = "Fundamentals of Life Sciences\nDr. Elena Rodriguez\nSpring 2025"

    # --- Slide 2: Course Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Course Overview"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Cell Biology and Molecular Structures"
    tf2.add_paragraph().text = "Genetics and Heredity"
    tf2.add_paragraph().text = "Evolutionary Biology"
    tf2.add_paragraph().text = "Ecology and Ecosystems"
    tf2.add_paragraph().text = "Human Physiology"

    # --- Slide 3: What is Biology? ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "What is Biology?"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Biology is the scientific study of life and living organisms"
    tf3.add_paragraph().text = "Encompasses molecular processes to ecosystem dynamics"
    tf3.add_paragraph().text = "Key branches: botany, zoology, microbiology, genetics"
    tf3.add_paragraph().text = "Interdisciplinary connections with chemistry and physics"

    # --- Slide 4: The Cell - Building Block of Life ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "The Cell - Building Block of Life"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "All living organisms are composed of one or more cells"
    tf4.add_paragraph().text = "Prokaryotic cells: no membrane-bound nucleus (bacteria, archaea)"
    tf4.add_paragraph().text = "Eukaryotic cells: membrane-bound organelles (plants, animals, fungi)"
    tf4.add_paragraph().text = "Cell theory established by Schleiden and Schwann (1838-1839)"

    # --- Slide 5: DNA and Genetic Information ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "DNA and Genetic Information"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "DNA (deoxyribonucleic acid) is the hereditary material in all organisms"
    tf5.add_paragraph().text = "Double helix structure discovered by Watson and Crick (1953)"
    tf5.add_paragraph().text = "Contains four nucleotide bases: Adenine, Thymine, Guanine, Cytosine"
    tf5.add_paragraph().text = "Genes are segments of DNA that encode proteins"

    # --- Slide 6: Gene Expression ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Gene Expression"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Transcription: DNA is copied into messenger RNA (mRNA)"
    tf6.add_paragraph().text = "Translation: mRNA is decoded by ribosomes to build proteins"
    tf6.add_paragraph().text = "Central dogma: DNA -> RNA -> Protein"
    tf6.add_paragraph().text = "Epigenetic modifications can regulate gene activity"

    # --- Slide 7: Evolution and Natural Selection ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Evolution and Natural Selection"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "Charles Darwin proposed natural selection as a mechanism of evolution"
    tf7.add_paragraph().text = "Organisms with favorable traits are more likely to survive and reproduce"
    tf7.add_paragraph().text = "Genetic variation arises through mutation and recombination"
    tf7.add_paragraph().text = "Over time, populations adapt to their environments"

    # --- Slide 8: Biodiversity ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Biodiversity"
    tf8 = slide8.placeholders[1].text_frame
    tf8.text = "Earth is home to an estimated 8.7 million species"
    tf8.add_paragraph().text = "Three levels: genetic diversity, species diversity, ecosystem diversity"
    tf8.add_paragraph().text = "Biodiversity hotspots contain high concentrations of endemic species"
    tf8.add_paragraph().text = "Human activities threaten biodiversity through habitat loss and climate change"

    # --- Slide 9: Ecosystems and Energy Flow ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Ecosystems and Energy Flow"
    tf9 = slide9.placeholders[1].text_frame
    tf9.text = "Energy flows through ecosystems via food chains and food webs"
    tf9.add_paragraph().text = "Producers (autotrophs) convert sunlight to chemical energy"
    tf9.add_paragraph().text = "Consumers (heterotrophs) obtain energy by feeding on other organisms"
    tf9.add_paragraph().text = "Decomposers recycle nutrients back into the ecosystem"

    # --- Slide 10: Key Takeaways (EMPTY - task target) ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    # Add just the title text box
    txBox = slide10.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf10 = txBox.text_frame
    p = tf10.paragraphs[0]
    p.text = "Key Takeaways"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
