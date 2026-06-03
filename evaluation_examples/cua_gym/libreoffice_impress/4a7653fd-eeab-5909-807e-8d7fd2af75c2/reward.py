"""
Reward Script: Create lecture summary slide with three takeaway boxes
Task ID: impress_teach_067
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Three rounded rectangles exist on slide 10
  Component 2 (0.30): Correct fill colors (#E3F2FD, #F3E5F5, #E8F5E9)
  Component 3 (0.25): Correct text in each box
  Component 4 (0.15): Equal vertical spacing between boxes
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_067'

# Expected values from task description
EXPECTED_FILLS = ['E3F2FD', 'F3E5F5', 'E8F5E9']
EXPECTED_TEXTS = [
    'Takeaway 1: Cells are the basic unit of life',
    'Takeaway 2: DNA carries genetic information',
    'Takeaway 3: Evolution drives biodiversity',
]


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 10 slides
    if len(prs.slides) < 10:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 10")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[9]  # slide 10 (0-indexed)

    # Find all rounded rectangle auto shapes on slide 10
    rounded_rects = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                from pptx.enum.shapes import MSO_SHAPE
                if shape.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
                    rounded_rects.append(shape)
            except Exception:
                # Fallback: check shape name or XML
                if 'round' in shape.name.lower():
                    rounded_rects.append(shape)

    # Sort by vertical position (top) to get order
    rounded_rects.sort(key=lambda s: s.top)

    # Component 1: Three rounded rectangles exist on slide 10 (0.30 points)
    try:
        if len(rounded_rects) == 3:
            print(f"PASS: Component 1 — Found 3 rounded rectangles on slide 10 (0.30 pts)")
            total_score += 0.30
        elif len(rounded_rects) >= 1:
            partial = 0.10 * min(len(rounded_rects), 3)
            print(f"PARTIAL: Component 1 — Found {len(rounded_rects)} rounded rectangles, expected 3 ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — No rounded rectangles found on slide 10")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Correct fill colors (0.30 points — 0.10 per box)
    for idx, expected_color in enumerate(EXPECTED_FILLS):
        try:
            if idx < len(rounded_rects):
                shape = rounded_rects[idx]
                fill = shape.fill
                if fill.type is not None and fill.type == 1:  # solid fill
                    actual_color = str(fill.fore_color.rgb)
                    if actual_color.upper() == expected_color.upper():
                        print(f"PASS: Component 2.{idx+1} — Box {idx+1} fill color {actual_color} matches expected {expected_color} (0.10 pts)")
                        total_score += 0.10
                    else:
                        print(f"FAIL: Component 2.{idx+1} — Box {idx+1} fill color {actual_color}, expected {expected_color}")
                else:
                    print(f"FAIL: Component 2.{idx+1} — Box {idx+1} has no solid fill (type={fill.type})")
            else:
                print(f"FAIL: Component 2.{idx+1} — Box {idx+1} not found")
        except Exception as e:
            print(f"ERROR: Component 2.{idx+1} — Could not read fill: {e}")

    # Component 3: Correct text in each box (0.25 points — ~0.083 per box)
    pts_per_box = 0.25 / 3.0
    for idx, expected_text in enumerate(EXPECTED_TEXTS):
        try:
            if idx < len(rounded_rects):
                shape = rounded_rects[idx]
                actual_text = shape.text.strip() if hasattr(shape, 'text') else ''
                if actual_text == expected_text:
                    print(f"PASS: Component 3.{idx+1} — Box {idx+1} text matches ({pts_per_box:.3f} pts)")
                    total_score += pts_per_box
                elif expected_text.lower() in actual_text.lower():
                    print(f"PARTIAL: Component 3.{idx+1} — Box {idx+1} text contains expected (partial credit)")
                    total_score += pts_per_box * 0.5
                else:
                    print(f"FAIL: Component 3.{idx+1} — Box {idx+1} text={repr(actual_text[:80])}, expected={repr(expected_text[:80])}")
            else:
                print(f"FAIL: Component 3.{idx+1} — Box {idx+1} not found")
        except Exception as e:
            print(f"ERROR: Component 3.{idx+1} — {e}")

    # Component 4: Equal vertical spacing between boxes (0.15 points)
    try:
        if len(rounded_rects) == 3:
            box1, box2, box3 = rounded_rects
            gap1 = box2.top - (box1.top + box1.height)
            gap2 = box3.top - (box2.top + box2.height)
            # Allow 5% tolerance on spacing equality
            if gap1 > 0 and gap2 > 0:
                max_gap = max(abs(gap1), abs(gap2))
                if max_gap > 0 and abs(gap1 - gap2) / max_gap <= 0.05:
                    print(f"PASS: Component 4 — Equal spacing: gap1={gap1}, gap2={gap2} (0.15 pts)")
                    total_score += 0.15
                else:
                    print(f"FAIL: Component 4 — Unequal spacing: gap1={gap1}, gap2={gap2}")
            else:
                print(f"FAIL: Component 4 — Boxes overlap or touch: gap1={gap1}, gap2={gap2}")
        else:
            print(f"FAIL: Component 4 — Need exactly 3 boxes to check spacing, found {len(rounded_rects)}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.4f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
