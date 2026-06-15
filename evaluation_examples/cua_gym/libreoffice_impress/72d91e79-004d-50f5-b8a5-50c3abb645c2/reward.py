"""
FINAL REWARD SCRIPT - SUCCESS
Task: While tidying up a 90-slide deck in LibreOffice Impress, I spotted that the heading on slide 79 is still in Liberation Sans. What steps do I follow to change that title to DejaVu Serif, 44 pt, Bold so it matches the rest of the presentation?
Generated: 2025-09-10 23:32:39
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation
from pptx.util import Pt

"""
Reward Script: Verify that the title on slide 79 is set to DejaVu Serif, 44 pt, Bold

Scoring Breakdown (progressive, out of 1.0):
• 0.1  – Target slide exists in the deck
• 0.1  – Target slide contains at least one text run to inspect
• 0.25 – Font family on any run matches DejaVu Serif (case-insensitive)
• 0.25 – Bold property on any run is True
• 0.30 – Font size on any run is 44 pt (±1 pt tolerance)

A perfect match yields a score of 1.0.  Partial matches yield proportionally
lower scores.
"""

FILE_PATH = (
    "/home/user/while_tidying_up_a_90_slide_deck_in_libreoffice_impress_i_spotted_"
    "that_the_heading_on_slide_79_is_st_golden.pptx"
)
TARGET_SLIDE_NUMBER = 79  # 1-based slide index from the task description
EXPECTED_FONT_NAME = "DejaVu Serif"
EXPECTED_FONT_SIZE_PT = 44
EXPECTED_BOLD = True


def emu_to_pt(emu_val):
    """Convert EMU value from python-pptx to points (Pt)."""
    if emu_val is None:
        return None
    # 1 point  = 12700 EMU in OOXML
    return emu_val / 12700


def collect_runs(slide):
    """Return a list of dicts with font information for every text run on slide."""
    runs_info = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                font = run.font
                runs_info.append(
                    {
                        "text": run.text,
                        "font_name": font.name,
                        "size_pt": emu_to_pt(font.size),
                        "bold": font.bold,
                    }
                )
    return runs_info


def evaluate_presentation(path):
    total_score = 0.0
    max_score = 1.0

    # --- prerequisite: file must exist and load (0 pts) ---
    if not os.path.exists(path):
        print("✗ File not found – task failed.")
        return 0.0

    try:
        prs = Presentation(path)
        print(f"✓ Loaded presentation with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Could not open presentation: {e}")
        return 0.0

    # --- slide existence (0.1) ---
    if len(prs.slides) < TARGET_SLIDE_NUMBER:
        print(
            f"✗ Slide {TARGET_SLIDE_NUMBER} does not exist – only {len(prs.slides)} slides found."
        )
        return total_score
    print(f"✓ Slide {TARGET_SLIDE_NUMBER} exists")
    total_score += 0.1

    # --- gather text runs (0.1 if any found) ---
    slide = prs.slides[TARGET_SLIDE_NUMBER - 1]  # zero-based index
    runs = collect_runs(slide)
    if not runs:
        print("✗ No text found on the target slide")
        return total_score
    print(f"✓ Located {len(runs)} text run(s) on slide {TARGET_SLIDE_NUMBER}")
    total_score += 0.1

    # --- evaluate font attributes on any run ---
    name_ok = False
    size_ok = False
    bold_ok = False

    for info in runs:
        preview = info["text"].replace("\n", " ")[:30]
        print(
            f"  • '{preview}' | Font='{info['font_name']}' | "
            f"Size={info['size_pt']} pt | Bold={info['bold']}"
        )

        if info["font_name"] and EXPECTED_FONT_NAME.lower() in info["font_name"].lower():
            name_ok = True
        if info["bold"] is EXPECTED_BOLD:
            bold_ok = True
        if (
            info["size_pt"] is not None
            and abs(info["size_pt"] - EXPECTED_FONT_SIZE_PT) <= 1
        ):
            size_ok = True

    # --- scoring ---
    if name_ok:
        print("✓ Font family matches DejaVu Serif")
        total_score += 0.25
    else:
        print("✗ Font family mismatch")

    if bold_ok:
        print("✓ Bold property is set correctly")
        total_score += 0.25
    else:
        print("✗ Bold property mismatch")

    if size_ok:
        print("✓ Font size is 44 pt (±1 pt)")
        total_score += 0.30
    else:
        print("✗ Font size mismatch")

    final = min(total_score, max_score)
    print(f"Total score: {final} / {max_score}")
    return final


if __name__ == "__main__":
    reward = evaluate_presentation(FILE_PATH)
    print(f"REWARD: {reward}")

