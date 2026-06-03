"""
FINAL REWARD SCRIPT - SUCCESS
Task: I need the title on the first slide to really stand out, so can you switch it all to uppercase letters?
Generated: 2025-08-07 10:55:02
Status: success
Model: o4-mini
Total Steps: 1
"""

import os
from pptx import Presentation


def verify_uppercase_title(file_path):
    """
    Verifies that the title on the first slide of the PowerPoint presentation
    is entirely uppercase. Uses progressive scoring:
      - File existence check: 0.2
      - Presentation load: 0.1
      - At least one slide present: 0.1
      - Title shape exists with non-empty text: 0.3
      - Title text is all uppercase: 0.3
    Returns the final score between 0.0 and 1.0 and prints detailed steps.
    """
    print("Checking task completion: Uppercase title on first slide")
    total_score = 0.0
    max_score = 1.0

    # Requirement 1: File existence (0.2 points)
    try:
        if os.path.isfile(file_path):
            print("✓ File exists (0.2 points)")
            total_score += 0.2
        else:
            print(f"✗ File not found: {file_path} (0 points)")
            print(f"REWARD: {total_score}")
            return total_score
    except Exception as e:
        print(f"✗ Error checking file existence: {e} (0 points)")
        print(f"REWARD: {total_score}")
        return total_score

    # Requirement 2: Load presentation (0.1 points)
    try:
        prs = Presentation(file_path)
        print("✓ Presentation loaded (0.1 points)")
        total_score += 0.1
    except Exception as e:
        print(f"✗ Error loading presentation: {e} (0 points)")
        print(f"REWARD: {total_score}")
        return total_score

    # Requirement 3: At least one slide present (0.1 points)
    slide_count = len(prs.slides)
    print(f"Slide count found: {slide_count}")
    if slide_count >= 1:
        print("✓ At least one slide present (0.1 points)")
        total_score += 0.1
    else:
        print("✗ No slides in presentation (0 points)")
        print(f"REWARD: {total_score}")
        return total_score

    # Requirement 4: Title shape exists with non-empty text on first slide (0.3 points)
    first_slide = prs.slides[0]
    title_shape = first_slide.shapes.title
    if title_shape and hasattr(title_shape, "text"):
        title_text = title_shape.text.strip()
        if title_text:
            print(f"✓ Title found on first slide: '{title_text}' (0.3 points)")
            total_score += 0.3
        else:
            print("✗ Title shape found but text is empty (0 points)")
            print(f"REWARD: {total_score}")
            return total_score
    else:
        print("✗ Title shape not found on first slide (0 points)")
        print(f"REWARD: {total_score}")
        return total_score

    # Requirement 5: Title text is all uppercase (0.3 points)
    if title_text == title_text.upper():
        print("✓ Title text is all uppercase (0.3 points)")
        total_score += 0.3
    else:
        print(f"✗ Title text is not uppercase: '{title_text}' (0 points)")

    # Final score calculation
    final_score = min(total_score, max_score)
    print(f"Final score: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score


if __name__ == '__main__':
    # Path to the presentation to verify
    file_path = '/home/user/i_need_the_title_on_the_first_slide_to_really_stand_out_so_can_you_switch_it_all_to_uppercase_letter.pptx'
    verify_uppercase_title(file_path)

