"""
Initial Setup: Table_Cleanup presentation with inconsistently formatted tables
Task ID: impress_exec_094
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_094'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_cell_fill(cell, hex_color):
    """Set solid fill on a table cell."""
    tcPr = cell._tc.get_or_add_tcPr()
    # Remove existing solidFill
    for sf in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(sf)
    solidFill = tcPr.makeelement(qn('a:solidFill'), {})
    srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': hex_color})
    solidFill.append(srgbClr)
    tcPr.append(solidFill)


def set_cell_borders(cell, width_emu, hex_color):
    """Set borders on a table cell via XML."""
    tcPr = cell._tc.get_or_add_tcPr()
    # Remove existing borders
    for old in tcPr.findall(qn('a:tcBorders')):
        tcPr.remove(old)
    tcBorders = tcPr.makeelement(qn('a:tcBorders'), {})
    for side in ['lnL', 'lnR', 'lnT', 'lnB']:
        ln = tcBorders.makeelement(qn('a:' + side), {'w': str(width_emu), 'cmpd': 'sng'})
        sf = ln.makeelement(qn('a:solidFill'), {})
        clr = sf.makeelement(qn('a:srgbClr'), {'val': hex_color})
        sf.append(clr)
        ln.append(sf)
        tcBorders.append(ln)
    tcPr.append(tcBorders)


def format_cell_text(cell, font_name, font_size_pt, bold, color_rgb):
    """Format all text in a cell."""
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size_pt)
            run.font.bold = bold
            run.font.color.rgb = color_rgb


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q4 2025 Business Review"
    slide1.placeholders[1].text = "Prepared by Sarah Chen, VP of Operations"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    tf = slide2.placeholders[1].text_frame
    tf.text = "Regional Sales Performance"
    tf.add_paragraph().text = "Product Line Analysis"
    tf.add_paragraph().text = "Customer Satisfaction Metrics"
    tf.add_paragraph().text = "Key Initiatives for Q1 2026"
    tf.add_paragraph().text = "Budget Allocation Summary"

    # --- Slide 3: Regional Sales Table (INCONSISTENT formatting #1) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Regional Sales Performance - Q4 2025"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    headers3 = ["Region", "Revenue ($K)", "Target ($K)", "% Achieved", "Growth YoY"]
    data3 = [
        ["North America", "$2,450", "$2,300", "106.5%", "+12.3%"],
        ["Europe", "$1,890", "$2,000", "94.5%", "+8.7%"],
        ["Asia Pacific", "$1,670", "$1,500", "111.3%", "+18.2%"],
        ["Latin America", "$890", "$950", "93.7%", "+5.1%"],
        ["Middle East & Africa", "$520", "$480", "108.3%", "+22.4%"],
    ]
    rows = len(data3) + 1
    cols = len(headers3)
    tbl_shape3 = slide3.shapes.add_table(rows, cols, Inches(0.5), Inches(1.3), Inches(11), Inches(3))
    tbl3 = tbl_shape3.table

    # Inconsistent style #1: red header, large comic-sans-like font, thick black borders
    for c, h in enumerate(headers3):
        cell = tbl3.cell(0, c)
        cell.text = h
        set_cell_fill(cell, 'CC0000')  # red
        format_cell_text(cell, 'Arial Black', 16, True, RGBColor(0xFF, 0xFF, 0x00))  # yellow text
        set_cell_borders(cell, 19050, '000000')  # 1.5pt black

    for r, row_data in enumerate(data3, 1):
        for c, val in enumerate(row_data):
            cell = tbl3.cell(r, c)
            cell.text = val
            set_cell_fill(cell, 'FFFFCC' if r % 2 == 0 else 'FFE0E0')  # yellow/pink alternating
            format_cell_text(cell, 'Times New Roman', 11, False, RGBColor(0x00, 0x00, 0x00))
            set_cell_borders(cell, 19050, '000000')  # 1.5pt black

    # --- Slide 4: Revenue Chart placeholder ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Revenue Trend - Last 4 Quarters"
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    txBox2 = slide4.shapes.add_textbox(Inches(2), Inches(2.5), Inches(8), Inches(2))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "[Revenue chart visualization placeholder - Q1-Q4 2025 trend data]"
    run2 = p2.runs[0]
    run2.font.size = Pt(14)
    run2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # --- Slide 5: Product Line Table (INCONSISTENT formatting #2) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Product Line Analysis - Q4 2025"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    headers5 = ["Product", "Units Sold", "Revenue ($K)", "Margin %", "Status"]
    data5 = [
        ["Enterprise Suite Pro", "1,245", "$3,120", "42.5%", "Growing"],
        ["CloudSync Platform", "3,890", "$1,945", "58.2%", "Stable"],
        ["DataVault Express", "890", "$712", "35.8%", "Declining"],
        ["SecureNet Gateway", "2,100", "$2,520", "51.3%", "Growing"],
        ["AnalyticsHub Core", "1,567", "$1,254", "47.1%", "Stable"],
        ["MobileFirst SDK", "4,230", "$845", "62.7%", "Growing"],
    ]
    rows = len(data5) + 1
    cols = len(headers5)
    tbl_shape5 = slide5.shapes.add_table(rows, cols, Inches(0.5), Inches(1.3), Inches(11), Inches(3.5))
    tbl5 = tbl_shape5.table

    # Inconsistent style #2: dark green header, small font, no borders visible
    for c, h in enumerate(headers5):
        cell = tbl5.cell(0, c)
        cell.text = h
        set_cell_fill(cell, '006600')  # dark green
        format_cell_text(cell, 'Courier New', 10, False, RGBColor(0xCC, 0xFF, 0xCC))  # light green text, not bold
        set_cell_borders(cell, 3175, 'FFFFFF')  # 0.25pt white (invisible)

    for r, row_data in enumerate(data5, 1):
        for c, val in enumerate(row_data):
            cell = tbl5.cell(r, c)
            cell.text = val
            set_cell_fill(cell, 'E8F5E9' if r % 2 == 0 else 'FFFFFF')  # green-tinted alternating
            format_cell_text(cell, 'Verdana', 10, False, RGBColor(0x00, 0x66, 0x00))  # dark green text
            set_cell_borders(cell, 3175, 'FFFFFF')  # barely visible

    # --- Slide 6: Key Takeaways ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Key Takeaways"
    tf = slide6.placeholders[1].text_frame
    tf.text = "North America exceeded targets by 6.5%, driven by Enterprise Suite adoption"
    tf.add_paragraph().text = "Asia Pacific showed strongest growth at 18.2% YoY"
    tf.add_paragraph().text = "CloudSync Platform maintains highest margins at 58.2%"
    tf.add_paragraph().text = "MobileFirst SDK fastest growing product with 4,230 units"

    # --- Slide 7: Customer Insights ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Customer Satisfaction Overview"
    tf = slide7.placeholders[1].text_frame
    tf.text = "NPS Score: 72 (up from 65 in Q3)"
    tf.add_paragraph().text = "Customer retention rate: 94.2%"
    tf.add_paragraph().text = "Support ticket resolution: avg 4.2 hours"
    tf.add_paragraph().text = "Top request: improved mobile experience"

    # --- Slide 8: Customer Satisfaction Table (INCONSISTENT formatting #3) ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide8.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Customer Satisfaction Metrics by Channel"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    headers8 = ["Channel", "NPS Score", "CSAT %", "Avg Resolution (hrs)", "Tickets/Month"]
    data8 = [
        ["Phone Support", "78", "91.2%", "2.1", "3,450"],
        ["Email Support", "65", "84.5%", "8.4", "5,230"],
        ["Live Chat", "82", "93.7%", "1.3", "7,890"],
        ["Self-Service Portal", "71", "87.1%", "N/A", "12,100"],
        ["Social Media", "59", "76.3%", "6.2", "1,870"],
    ]
    rows = len(data8) + 1
    cols = len(headers8)
    tbl_shape8 = slide8.shapes.add_table(rows, cols, Inches(0.5), Inches(1.3), Inches(11), Inches(3))
    tbl8 = tbl_shape8.table

    # Inconsistent style #3: bright blue header, mixed font sizes, orange borders
    for c, h in enumerate(headers8):
        cell = tbl8.cell(0, c)
        cell.text = h
        set_cell_fill(cell, '0066FF')  # bright blue
        format_cell_text(cell, 'Georgia', 13, True, RGBColor(0xFF, 0xFF, 0xFF))
        set_cell_borders(cell, 12700, 'FF6600')  # 1pt orange

    for r, row_data in enumerate(data8, 1):
        for c, val in enumerate(row_data):
            cell = tbl8.cell(r, c)
            cell.text = val
            set_cell_fill(cell, 'CCE5FF' if r % 2 == 0 else 'FFFFFF')  # blue-tinted alternating
            format_cell_text(cell, 'Arial', 11, False, RGBColor(0x33, 0x33, 0x99))  # blue-ish text
            set_cell_borders(cell, 12700, 'FF6600')  # 1pt orange

    # --- Slide 9: Initiatives ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Q1 2026 Key Initiatives"
    tf = slide9.placeholders[1].text_frame
    tf.text = "Launch Enterprise Suite Pro v3.0 with AI capabilities"
    tf.add_paragraph().text = "Expand APAC sales team by 25%"
    tf.add_paragraph().text = "Migrate DataVault Express to cloud-native architecture"
    tf.add_paragraph().text = "Implement unified customer feedback platform"
    tf.add_paragraph().text = "Reduce support resolution time to under 3 hours"

    # --- Slide 10: Thank You ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[0])
    slide10.shapes.title.text = "Thank You"
    slide10.placeholders[1].text = "Questions? Contact: sarah.chen@company.com"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
