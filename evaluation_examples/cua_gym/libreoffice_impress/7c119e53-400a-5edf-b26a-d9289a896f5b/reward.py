"""
Reward Script: Set all tables in the presentation to have consistent formatting
Task ID: impress_exec_094
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): Header row fill #003366 on slides 3, 5, 8
  Component 2 (0.20): Header text white, bold, 14pt on slides 3, 5, 8
  Component 3 (0.20): Data rows alternating white/#F5F7FA fills on slides 3, 5, 8
  Component 4 (0.20): Data text 12pt Calibri #333333 on slides 3, 5, 8
  Component 5 (0.20): All borders 0.5pt #D0D0D0 on slides 3, 5, 8
"""

import os
import zipfile
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_094'
FILE_PATH = os.path.join(WORKDIR, f'{TASK_ID}.pptx')

# Target table slides (0-indexed)
TABLE_SLIDES = [2, 4, 7]  # slides 3, 5, 8

# Expected formatting values
HEADER_FILL = '003366'
HEADER_TEXT_COLOR = 'FFFFFF'
HEADER_BOLD = True
HEADER_SIZE = Pt(14)  # 177800 EMU
DATA_FONT_NAME = 'Calibri'
DATA_TEXT_COLOR = '333333'
DATA_SIZE = Pt(12)  # 152400 EMU
ALT_FILLS = ['FFFFFF', 'F5F7FA']  # alternating row fills for data rows (odd=white, even=F5F7FA)

# Border expected values
BORDER_WIDTH = '6350'  # 0.5pt in EMU
BORDER_COLOR = 'D0D0D0'


def get_cell_fill_rgb(cell):
    """Get the fill color of a table cell as uppercase hex string, or None."""
    try:
        fill = cell.fill
        if fill.type is not None:
            return str(fill.fore_color.rgb).upper()
    except Exception:
        pass
    return None


def get_run_color_rgb(run):
    """Get font color of a run as uppercase hex string, or None."""
    try:
        if run.font.color.type is not None:
            return str(run.font.color.rgb).upper()
    except Exception:
        pass
    return None


def check_borders_via_xml(pptx_path, slide_indices):
    """
    Check that all table cells on given slides have 4 borders with correct width and color.
    Returns (pass_count, total_count) of border sets checked.
    """
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    border_names = [f'{{{ns_a}}}lnL', f'{{{ns_a}}}lnR', f'{{{ns_a}}}lnT', f'{{{ns_a}}}lnB']

    total_cells = 0
    pass_cells = 0

    with zipfile.ZipFile(pptx_path, 'r') as zf:
        for si in slide_indices:
            slide_file = f'ppt/slides/slide{si + 1}.xml'
            try:
                with zf.open(slide_file) as f:
                    root = ET.parse(f).getroot()
            except KeyError:
                continue

            # Find all tc elements (table cells)
            for tc in root.iter(f'{{{ns_a}}}tc'):
                tcPr = tc.find(f'{{{ns_a}}}tcPr')
                if tcPr is None:
                    total_cells += 1
                    continue

                total_cells += 1
                # Borders are nested inside tcBorders element
                tcBorders = tcPr.find(f'{{{ns_a}}}tcBorders')
                if tcBorders is None:
                    # Also check directly under tcPr as fallback
                    tcBorders = tcPr

                cell_ok = True
                for bn in border_names:
                    ln = tcBorders.find(bn)
                    if ln is None:
                        cell_ok = False
                        break
                    w = ln.get('w', '')
                    if w != BORDER_WIDTH:
                        cell_ok = False
                        break
                    sf = ln.find(f'{{{ns_a}}}solidFill')
                    if sf is None:
                        cell_ok = False
                        break
                    clr = sf.find(f'{{{ns_a}}}srgbClr')
                    if clr is None or clr.get('val', '').upper() != BORDER_COLOR:
                        cell_ok = False
                        break

                if cell_ok:
                    pass_cells += 1

    return pass_cells, total_cells


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)
    if len(slides) < 8:
        print(f"CRITICAL: Expected at least 8 slides, found {len(slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Collect tables from target slides
    tables_info = []  # list of (slide_index, table) tuples
    for si in TABLE_SLIDES:
        slide = slides[si]
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                tables_info.append((si, shape.table))
                break

    if len(tables_info) != 3:
        print(f"CRITICAL: Expected 3 tables (on slides 3, 5, 8), found {len(tables_info)}")
        print("REWARD: 0.0")
        return 0.0

    # ---- Component 1: Header row fill #003366 (0.20 points) ----
    try:
        comp1_pass = 0
        comp1_total = 0
        for si, table in tables_info:
            cols = len(table.columns)
            for ci in range(cols):
                comp1_total += 1
                fill_rgb = get_cell_fill_rgb(table.cell(0, ci))
                if fill_rgb == HEADER_FILL.upper():
                    comp1_pass += 1

        if comp1_total > 0 and comp1_pass == comp1_total:
            print(f"PASS: Component 1 -- Header fill #003366 on all {comp1_total} header cells (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 -- Header fill: {comp1_pass}/{comp1_total} cells have #003366")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # ---- Component 2: Header text white, bold, 14pt (0.20 points) ----
    try:
        comp2_pass = 0
        comp2_total = 0
        for si, table in tables_info:
            cols = len(table.columns)
            for ci in range(cols):
                cell = table.cell(0, ci)
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if not (run.text or "").strip():
                            continue
                        comp2_total += 1
                        color_ok = get_run_color_rgb(run) == HEADER_TEXT_COLOR
                        bold_ok = run.font.bold is True
                        size_ok = run.font.size is not None and run.font.size == HEADER_SIZE
                        if color_ok and bold_ok and size_ok:
                            comp2_pass += 1
                        else:
                            details = []
                            if not color_ok:
                                details.append(f"color={get_run_color_rgb(run)}")
                            if not bold_ok:
                                details.append(f"bold={run.font.bold}")
                            if not size_ok:
                                details.append(f"size={run.font.size}")
                            print(f"  DETAIL: Slide {si+1} header[0,{ci}] run '{run.text}': {', '.join(details)}")

        if comp2_total > 0 and comp2_pass == comp2_total:
            print(f"PASS: Component 2 -- Header text white/bold/14pt on all {comp2_total} runs (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 2 -- Header text: {comp2_pass}/{comp2_total} runs correct")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # ---- Component 3: Data rows alternating white/#F5F7FA fills (0.20 points) ----
    try:
        comp3_pass = 0
        comp3_total = 0
        for si, table in tables_info:
            rows = len(table.rows)
            cols = len(table.columns)
            for ri in range(1, rows):  # skip header
                expected_fill = ALT_FILLS[(ri - 1) % 2].upper()  # row1=white, row2=F5F7FA, ...
                for ci in range(cols):
                    comp3_total += 1
                    fill_rgb = get_cell_fill_rgb(table.cell(ri, ci))
                    if fill_rgb == expected_fill:
                        comp3_pass += 1
                    else:
                        print(f"  DETAIL: Slide {si+1} data[{ri},{ci}] fill={fill_rgb}, expected={expected_fill}")

        if comp3_total > 0 and comp3_pass == comp3_total:
            print(f"PASS: Component 3 -- Alternating data fills on all {comp3_total} data cells (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 3 -- Alternating fills: {comp3_pass}/{comp3_total} cells correct")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # ---- Component 4: Data text 12pt Calibri #333333 (0.20 points) ----
    try:
        comp4_pass = 0
        comp4_total = 0
        for si, table in tables_info:
            rows = len(table.rows)
            cols = len(table.columns)
            for ri in range(1, rows):
                for ci in range(cols):
                    cell = table.cell(ri, ci)
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            if not (run.text or "").strip():
                                continue
                            comp4_total += 1
                            font_ok = run.font.name == DATA_FONT_NAME
                            size_ok = run.font.size is not None and run.font.size == DATA_SIZE
                            color_ok = get_run_color_rgb(run) == DATA_TEXT_COLOR.upper()
                            if font_ok and size_ok and color_ok:
                                comp4_pass += 1
                            else:
                                details = []
                                if not font_ok:
                                    details.append(f"font={run.font.name}")
                                if not size_ok:
                                    details.append(f"size={run.font.size}")
                                if not color_ok:
                                    details.append(f"color={get_run_color_rgb(run)}")
                                print(f"  DETAIL: Slide {si+1} data[{ri},{ci}] run '{run.text}': {', '.join(details)}")

        if comp4_total > 0 and comp4_pass == comp4_total:
            print(f"PASS: Component 4 -- Data text 12pt Calibri #333333 on all {comp4_total} runs (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 4 -- Data text: {comp4_pass}/{comp4_total} runs correct")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # ---- Component 5: Borders 0.5pt #D0D0D0 on all table cells (0.20 points) ----
    try:
        pass_cells, total_cells = check_borders_via_xml(file_path, TABLE_SLIDES)
        if total_cells > 0 and pass_cells == total_cells:
            print(f"PASS: Component 5 -- Borders 0.5pt #D0D0D0 on all {total_cells} cells (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 5 -- Borders: {pass_cells}/{total_cells} cells have correct borders")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = min(total_score, 1.0)
    # Round to avoid floating point issues
    final_score = round(final_score, 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(FILE_PATH)
