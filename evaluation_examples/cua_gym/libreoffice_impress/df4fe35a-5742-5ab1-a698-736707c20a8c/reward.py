"""
Reward Script: Add note 'Introduce agenda here' to slide 1 and change its background to light blue.
Task ID: osworld_impress_note_bg_combined_001
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 1 notes contain 'Introduce agenda here'  — 0.5 points
  Component 2: Slide 1 background is light blue (ADD8E6)      — 0.5 points
  Total: 1.0
"""

import os
from pptx import Presentation
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_note_bg_combined_001'


def persist_app_state():
    """Send Ctrl+S to persist any unsaved GUI edits before verification."""
    try:
        os.environ["DISPLAY"] = ":0"
        import pyautogui
        import time
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Task: Add note 'Introduce agenda here' to slide 1 AND change slide 1 background to light blue.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 1 slide
    if len(prs.slides) < 1:
        print("CRITICAL: Presentation has no slides.")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]

    # Component 1: Slide 1 notes contain 'Introduce agenda here' (0.5 points)
    # Initial state: no notes (empty string). Golden state: 'Introduce agenda here'
    try:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        expected_note = "Introduce agenda here"
        if expected_note in notes_text:
            print(f"PASS: Component 1 — Slide 1 notes contain '{expected_note}' (found: {repr(notes_text)}) (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 1 — Expected notes to contain '{expected_note}', found: {repr(notes_text)}")
    except Exception as e:
        print(f"ERROR: Component 1 — Could not read slide 1 notes: {e}")

    # Component 2: Slide 1 background is light blue (ADD8E6) (0.5 points)
    # Initial state: white background (FFFFFF). Golden state: light blue (ADD8E6)
    EXPECTED_BG_RGB = "ADD8E6"
    try:
        fill = slide.background.fill
        actual_rgb = None
        if fill.type == 1:  # SOLID fill directly on slide
            actual_rgb = str(fill.fore_color.rgb).upper()
        elif fill.type == 5:  # Inherited from master — need to check master
            master_fill = slide.slide_layout.slide_master.background.fill
            if master_fill.type == 1:
                actual_rgb = str(master_fill.fore_color.rgb).upper()

        if actual_rgb == EXPECTED_BG_RGB:
            print(f"PASS: Component 2 — Slide 1 background is light blue (ADD8E6) (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 2 — Expected background ADD8E6, found: {repr(actual_rgb)}")
    except Exception as e:
        print(f"ERROR: Component 2 — Could not check slide 1 background: {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the VM environment
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
persist_app_state()
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
