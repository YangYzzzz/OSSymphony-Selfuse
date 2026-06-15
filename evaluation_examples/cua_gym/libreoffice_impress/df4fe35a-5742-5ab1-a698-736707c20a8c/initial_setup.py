"""
Initial Setup: Add note 'Introduce agenda here' to slide 1 and change its background to light blue.
Task ID: osworld_impress_note_bg_combined_001
Domain: libreoffice_impress

Creates a 5-slide kick-off meeting deck with:
  - Slide 1: white background, NO speaker notes (these are what the task will add)
  - Slides 2-5: various meeting content
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_note_bg_combined_001'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide (white background, NO notes) ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])

    # Explicitly set white background on slide 1
    fill1 = slide1.background.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    slide1.shapes.title.text = "Q2 Kick-Off Meeting"
    slide1.placeholders[1].text = "March 2025 | Strategy & Planning Session"

    # NOTE: No speaker notes on slide 1 — this is what the task requires the agent to add

    # ---- Slide 2: Agenda ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "1. Welcome & Introductions"
    items = [
        "2. Q1 Performance Review",
        "3. Q2 Goals & Objectives",
        "4. Team Project Updates",
        "5. Open Discussion & Q&A",
    ]
    for item in items:
        p = tf2.add_paragraph()
        p.text = item
        p.level = 0

    slide2.notes_slide.notes_text_frame.text = "Keep agenda tight — 5 minutes per section."

    # ---- Slide 3: Q1 Performance Review ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Q1 Performance Highlights"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Revenue: $4.2M (103% of target)"
    highlights = [
        "Customer Satisfaction: 92% (up 4pts)",
        "New Accounts: 47 (exceeded by 12%)",
        "Team Headcount: 38 (+6 new hires)",
        "Support Tickets Resolved: 1,240",
    ]
    for h in highlights:
        p = tf3.add_paragraph()
        p.text = h
        p.level = 0

    slide3.notes_slide.notes_text_frame.text = "Highlight the NPS improvement. Sarah Chen to present revenue section."

    # ---- Slide 4: Q2 Goals ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Q2 Objectives"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Drive 15% revenue growth quarter-over-quarter"
    goals = [
        "Launch 3 new product features by May 30",
        "Expand into 2 new regional markets",
        "Reduce customer churn below 5%",
        "Complete ISO 27001 certification",
    ]
    for g in goals:
        p = tf4.add_paragraph()
        p.text = g
        p.level = 0

    slide4.notes_slide.notes_text_frame.text = "Marcus Johnson to walk through regional expansion plan."

    # ---- Slide 5: Next Steps ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Next Steps & Action Items"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Finalize Q2 roadmap by April 5"
    actions = [
        "Schedule 1:1 check-ins with team leads",
        "Submit budget requests by April 10",
        "Distribute meeting summary within 24 hours",
        "Follow up with new client prospects",
    ]
    for a in actions:
        p = tf5.add_paragraph()
        p.text = a
        p.level = 0

    slide5.notes_slide.notes_text_frame.text = "Close with energy. Remind team of the off-site on April 18."

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup — open the file in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
