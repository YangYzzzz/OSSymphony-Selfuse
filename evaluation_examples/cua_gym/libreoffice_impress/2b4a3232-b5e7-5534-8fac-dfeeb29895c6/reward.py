"""
Reward Script: Multi-operation LibreOffice Impress task
Task ID: osworld_impress_multi_op_combined_009
Domain: libreoffice_impress

Task: Perform three updates:
  (1) Change the title on slide 1 to 'Q4 Business Review'
  (2) Underline and color the titles on slides 2 and 3 dark red
  (3) Remove the textbox on slide 5 that says 'Preliminary Data'

Scoring Rubric:
  Component 1 — Slide 1 title text changed to 'Q4 Business Review'  : 0.35 pts
  Component 2 — Slide 2 title underlined and colored dark red        : 0.30 pts
  Component 3 — Slide 3 title underlined and colored dark red        : 0.20 pts
  Component 4 — Slide 5 'Preliminary Data' textbox removed           : 0.15 pts
  Total: 1.0
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_multi_op_combined_009'

# Dark red color in hex — expected value from golden file exploration
DARK_RED_HEX = '8B0000'


def get_title_text(slide):
    """Return the text of the first title-like placeholder on a slide."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name.startswith('Title'):
            full_text = shape.text_frame.text
            return full_text
    return None


def get_title_shape(slide):
    """Return the title shape object from a slide."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name.startswith('Title'):
            return shape
    return None


def get_title_run_props(slide):
    """
    Return a dict with underline and color_hex for the title shape runs.
    We check all runs and collect their underline/color properties.
    Returns (underline, color_hex) tuple or (None, None) if no runs found.
    """
    title_shape = get_title_shape(slide)
    if title_shape is None:
        return None, None
    for para in title_shape.text_frame.paragraphs:
        runs = [r for r in para.runs if (r.text or '').strip()]
        if runs:
            run = runs[0]
            underline = run.font.underline
            color_hex = None
            try:
                if run.font.color.type is not None:
                    color_hex = str(run.font.color.rgb).upper()
            except Exception:
                color_hex = None
            return underline, color_hex
    return None, None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have 7 slides
    if len(prs.slides) != 7:
        print(f"CRITICAL: Expected 7 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide1 = prs.slides[0]
    slide2 = prs.slides[1]
    slide3 = prs.slides[2]
    slide5 = prs.slides[4]

    # Component 1: Slide 1 title text changed to 'Q4 Business Review' (0.35 points)
    try:
        title1_text = get_title_text(slide1)
        expected_title = 'Q4 Business Review'
        if title1_text and title1_text.strip() == expected_title:
            print(f"PASS: Component 1 — Slide 1 title is '{title1_text}' (0.35 pts)")
            total_score += 0.35
        else:
            print(f"FAIL: Component 1 — Slide 1 title expected '{expected_title}', found '{title1_text}'")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 2 title underlined and colored dark red (0.30 points)
    try:
        underline2, color2 = get_title_run_props(slide2)
        title2_text = get_title_text(slide2)
        dark_red = DARK_RED_HEX
        underline_ok = underline2 is True
        color_ok = color2 == dark_red if color2 is not None else False
        if underline_ok and color_ok:
            print(f"PASS: Component 2 — Slide 2 title '{title2_text}' is underlined and dark red #{color2} (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 2 — Slide 2 title '{title2_text}': underline={underline2} (expected True), color={color2} (expected {dark_red})")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 3 title underlined and colored dark red (0.20 points)
    try:
        underline3, color3 = get_title_run_props(slide3)
        title3_text = get_title_text(slide3)
        dark_red = DARK_RED_HEX
        underline_ok = underline3 is True
        color_ok = color3 == dark_red if color3 is not None else False
        if underline_ok and color_ok:
            print(f"PASS: Component 3 — Slide 3 title '{title3_text}' is underlined and dark red #{color3} (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 3 — Slide 3 title '{title3_text}': underline={underline3} (expected True), color={color3} (expected {dark_red})")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 5 'Preliminary Data' textbox is removed (0.15 points)
    try:
        preliminary_shapes = [
            shape for shape in slide5.shapes
            if shape.has_text_frame and 'Preliminary Data' in shape.text_frame.text.strip()
        ]
        if len(preliminary_shapes) == 0:
            print("PASS: Component 4 — Slide 5 'Preliminary Data' textbox is absent (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 — Slide 5 still contains {len(preliminary_shapes)} shape(s) with 'Preliminary Data'")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(round(total_score, 4), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point: test against canonical artifact path on this VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
