"""
Initial Setup: Create a 7-slide quarterly review presentation in initial state.
Task ID: osworld_impress_multi_op_combined_009
Domain: libreoffice_impress

Initial state:
- Slide 1 title reads 'Q3 Business Review' (NOT Q4)
- Slides 2 and 3 have plain black titles (no underline, no dark red)
- Slide 5 has a content textbox labeled 'Preliminary Data'
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_multi_op_combined_009'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Use standard 10x7.5 inch widescreen
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title slide — "Q3 Business Review" ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
    slide1.shapes.title.text = 'Q3 Business Review'
    # Subtitle
    subtitle = slide1.placeholders[1]
    subtitle.text = 'Annual Performance Summary\nFiscal Year 2024'

    # ---- Slide 2: Executive Summary ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide2.shapes.title.text = 'Executive Summary'
    # Set title to plain black — no underline, no special color
    title2_tf = slide2.shapes.title.text_frame
    for para in title2_tf.paragraphs:
        for run in para.runs:
            run.font.underline = False
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    content2 = slide2.placeholders[1]
    content2.text = (
        'Revenue grew 12% year-over-year\n'
        'Operating margin improved to 18.4%\n'
        'Customer retention rate: 94.2%\n'
        'New product launches: 3 major, 7 minor\n'
        'Headcount increased by 85 employees'
    )

    # ---- Slide 3: Financial Highlights ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide3.shapes.title.text = 'Financial Highlights'
    # Set title to plain black — no underline, no special color
    title3_tf = slide3.shapes.title.text_frame
    for para in title3_tf.paragraphs:
        for run in para.runs:
            run.font.underline = False
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    content3 = slide3.placeholders[1]
    content3.text = (
        'Total Revenue: $4.82B (+12% YoY)\n'
        'Gross Profit: $2.15B (44.6% margin)\n'
        'EBITDA: $887M (+9% YoY)\n'
        'Free Cash Flow: $612M\n'
        'Earnings Per Share: $3.47'
    )

    # ---- Slide 4: Regional Performance ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = 'Regional Performance'
    content4 = slide4.placeholders[1]
    content4.text = (
        'North America: $2.10B — up 8%\n'
        'Europe: $1.34B — up 15%\n'
        'Asia-Pacific: $0.98B — up 22%\n'
        'Latin America: $0.40B — up 6%\n'
        'Total International: 56% of revenue'
    )

    # ---- Slide 5: Market Outlook — with "Preliminary Data" textbox ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = 'Market Outlook'
    content5 = slide5.placeholders[1]
    content5.text = (
        'Industry growth forecast: 9% for Q4\n'
        'Emerging market opportunities in SEA\n'
        'Competitive positioning remains strong\n'
        'Three new partnerships signed in October'
    )
    # Add a textbox labeled "Preliminary Data" — this is what the agent must remove
    txbox = slide5.shapes.add_textbox(
        Inches(0.5), Inches(5.8), Inches(4.0), Inches(0.6)
    )
    tf = txbox.text_frame
    tf.text = 'Preliminary Data'
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.italic = True
    run.font.color.rgb = RGBColor(0xFF, 0x66, 0x00)  # orange-ish warning color

    # ---- Slide 6: Product Roadmap ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = 'Product Roadmap'
    content6 = slide6.placeholders[1]
    content6.text = (
        'Q4 2024: Launch of Enterprise Suite v3.0\n'
        'Q1 2025: AI-powered analytics module\n'
        'Q2 2025: Mobile platform overhaul\n'
        'Q3 2025: Global expansion of cloud services\n'
        'Ongoing: Customer feedback integration cycle'
    )

    # ---- Slide 7: Next Steps & Action Items ----
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = 'Next Steps & Action Items'
    content7 = slide7.placeholders[1]
    content7.text = (
        'Finalize Q4 budget allocations by Nov 15\n'
        'Complete hiring for 3 senior engineering roles\n'
        'Launch APAC partnership program in December\n'
        'Submit annual report draft by Dec 20\n'
        'Schedule all-hands meeting for January 2025'
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
