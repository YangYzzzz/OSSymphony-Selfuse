"""
Initial Setup: Create a blank presentation and open in LibreOffice Impress
Task ID: impress_wf_010
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_010'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    # Create a blank presentation with one blank slide
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only / blank-ish

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Ensure Desktop directory exists
    os.makedirs(f'{WORKDIR}/Desktop', exist_ok=True)

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
