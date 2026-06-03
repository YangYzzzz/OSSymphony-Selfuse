"""
Reward Script: Color Palette Showcase Presentation
Task ID: impress_wf_010
Domain: libreoffice_impress
Scoring:
  Component 1 (0.15) - File exists with exactly 4 slides
  Component 2 (0.30) - Correct background colors on all 4 slides
  Component 3 (0.30) - Correct title text (centered, white, bold, ~48pt)
  Component 4 (0.25) - Circle shape with hex code text on each slide
"""

import os
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_010'

# Expected slide data: (bg_color, title_text, hex_code)
EXPECTED_SLIDES = [
    ('E74C3C', 'Passion Red', '#E74C3C'),
    ('3498DB', 'Ocean Blue', '#3498DB'),
    ('2ECC71', 'Fresh Green', '#2ECC71'),
    ('F39C12', 'Warm Amber', '#F39C12'),
]


def get_all_text_shapes(slide):
    """Recursively get all shapes with text, including inside groups."""
    def extract(shape):
        results = []
        if hasattr(shape, "text") and hasattr(shape, "text_frame"):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.extend(extract(sub))
        return results
    out = []
    for shape in slide.shapes:
        out.extend(extract(shape))
    return out


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: File has exactly 4 slides (0.15 points)
    try:
        num_slides = len(prs.slides)
        if num_slides == 4:
            print(f"PASS: Component 1 -- File has exactly 4 slides (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 -- Expected 4 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Correct background colors on all 4 slides (0.30 points)
    # Each slide contributes 0.075 points
    try:
        bg_score = 0.0
        slides = list(prs.slides)
        for idx, (expected_bg, _, _) in enumerate(EXPECTED_SLIDES):
            if idx >= len(slides):
                print(f"FAIL: Component 2 -- Slide {idx+1} does not exist")
                continue
            slide = slides[idx]
            fill = slide.background.fill
            actual_bg = None
            if fill.type == 1:  # SOLID
                try:
                    actual_bg = str(fill.fore_color.rgb)
                except Exception:
                    pass
            elif fill.type == 5:  # inherited from master
                try:
                    master_fill = slide.slide_layout.slide_master.background.fill
                    if master_fill.type == 1:
                        actual_bg = str(master_fill.fore_color.rgb)
                except Exception:
                    pass

            if actual_bg and actual_bg.upper() == expected_bg.upper():
                print(f"PASS: Component 2 -- Slide {idx+1} bg color is {actual_bg} (correct)")
                bg_score += 0.075
            else:
                print(f"FAIL: Component 2 -- Slide {idx+1} bg expected {expected_bg}, found {actual_bg}")

        total_score += bg_score
        if bg_score >= 0.30:
            print(f"PASS: Component 2 -- All background colors correct (0.30 pts)")
        else:
            print(f"PARTIAL: Component 2 -- Background colors scored {bg_score:.3f}/0.30")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Correct title text - centered, white, bold, ~48pt (0.30 points)
    # Each slide contributes 0.075 points
    try:
        title_score = 0.0
        slides = list(prs.slides)
        for idx, (_, expected_title, _) in enumerate(EXPECTED_SLIDES):
            if idx >= len(slides):
                print(f"FAIL: Component 3 -- Slide {idx+1} does not exist")
                continue
            slide = slides[idx]
            text_shapes = get_all_text_shapes(slide)

            found_title = False
            for shape in text_shapes:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip().lower() == expected_title.lower():
                        # Found the title text, now check formatting
                        is_centered = para.alignment in (PP_ALIGN.CENTER, 2)
                        is_white = False
                        is_bold = False
                        size_ok = False
                        for run in para.runs:
                            if run.text.strip().lower() == expected_title.lower():
                                try:
                                    if run.font.color.type is not None:
                                        rgb_str = str(run.font.color.rgb).upper()
                                        is_white = rgb_str == 'FFFFFF'
                                except Exception:
                                    pass
                                is_bold = run.font.bold in (True,)
                                if run.font.size is not None:
                                    # 48pt = 609600 EMU. Allow tolerance: 44pt-52pt
                                    pt_size = run.font.size / 12700
                                    size_ok = 44 <= pt_size <= 52

                        if is_centered and is_white and is_bold and size_ok:
                            print(f"PASS: Component 3 -- Slide {idx+1} title '{expected_title}' formatted correctly")
                            title_score += 0.075
                            found_title = True
                            break
                        else:
                            details = f"centered={is_centered}, white={is_white}, bold={is_bold}, size_ok={size_ok}"
                            print(f"FAIL: Component 3 -- Slide {idx+1} title '{expected_title}' found but formatting wrong: {details}")
                            found_title = True
                            break
                if found_title:
                    break

            if not found_title:
                print(f"FAIL: Component 3 -- Slide {idx+1} title '{expected_title}' not found")

        total_score += title_score
        if title_score >= 0.30:
            print(f"PASS: Component 3 -- All titles correct (0.30 pts)")
        else:
            print(f"PARTIAL: Component 3 -- Titles scored {title_score:.3f}/0.30")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Circle/oval shape with hex code text on each slide (0.25 points)
    # Each slide contributes 0.0625 points
    try:
        hex_score = 0.0
        slides = list(prs.slides)
        for idx, (_, _, expected_hex) in enumerate(EXPECTED_SLIDES):
            if idx >= len(slides):
                print(f"FAIL: Component 4 -- Slide {idx+1} does not exist")
                continue
            slide = slides[idx]

            found_hex = False
            for shape in slide.shapes:
                # Check for AUTO_SHAPE (which includes ovals/circles) with text
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.has_text_frame:
                    shape_text = shape.text_frame.text.strip()
                    if expected_hex.upper() in shape_text.upper():
                        print(f"PASS: Component 4 -- Slide {idx+1} has shape with hex code '{shape_text}'")
                        hex_score += 0.0625
                        found_hex = True
                        break

            if not found_hex:
                # Also check group shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'shapes'):
                        for sub in shape.shapes:
                            if sub.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and sub.has_text_frame:
                                if expected_hex.upper() in sub.text_frame.text.strip().upper():
                                    print(f"PASS: Component 4 -- Slide {idx+1} has grouped shape with hex '{expected_hex}'")
                                    hex_score += 0.0625
                                    found_hex = True
                                    break
                        if found_hex:
                            break

            if not found_hex:
                print(f"FAIL: Component 4 -- Slide {idx+1} missing shape with hex code '{expected_hex}'")

        total_score += hex_score
        if hex_score >= 0.25:
            print(f"PASS: Component 4 -- All hex code shapes correct (0.25 pts)")
        else:
            print(f"PARTIAL: Component 4 -- Hex code shapes scored {hex_score:.4f}/0.25")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/Desktop/Color_Palette.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
