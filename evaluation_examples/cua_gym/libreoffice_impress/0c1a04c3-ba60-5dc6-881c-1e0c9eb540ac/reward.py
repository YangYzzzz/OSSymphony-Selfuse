"""
Reward Script: Create agenda slide with numbered list and custom formatting
Task ID: impress_rp_032
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): 6 agenda paragraphs exist on slide 2 with correct topic text
  Component 2 (0.25): Number runs are 28pt bold
  Component 3 (0.25): Number colors alternate teal/coral correctly
  Component 4 (0.25): Topic runs are 18pt and not bold
"""

import os
from pptx import Presentation
from pptx.util import Pt

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_032'

# Expected items: (number_str, topic_text, number_color_hex)
EXPECTED_ITEMS = [
    ('1', 'Welcome & Introductions', '008080'),
    ('2', 'Q1 Results Review', 'FF6F61'),
    ('3', 'Market Analysis', '008080'),
    ('4', 'Product Roadmap', 'FF6F61'),
    ('5', 'Q2 Goals', '008080'),
    ('6', 'Open Discussion', 'FF6F61'),
]


def find_agenda_textbox(slide):
    """Find the text shape on slide 2 that contains the numbered agenda items.
    We look for a shape with >= 6 paragraphs whose first run text is a digit."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        # Filter to non-empty paragraphs
        nonempty = [p for p in tf.paragraphs if p.text.strip()]
        if len(nonempty) >= 6:
            # Check if first non-empty paragraph starts with a digit run
            first_runs = [r for r in nonempty[0].runs if r.text.strip()]
            if first_runs and first_runs[0].text.strip().isdigit():
                return nonempty
    return None


def get_run_color_hex(run):
    """Safely get RGB hex string from a run's font color."""
    try:
        if run.font.color.type is not None:
            return str(run.font.color.rgb).upper()
    except Exception:
        pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 2")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[1]  # slide 2 (0-indexed)

    # Find the agenda text box with 6 items
    agenda_paras = find_agenda_textbox(slide)

    if agenda_paras is None:
        print("FAIL: No text shape found on slide 2 with 6+ numbered agenda items")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: 6 agenda paragraphs with correct topic text (0.25 points)
    try:
        matched_items = 0
        for i, (exp_num, exp_topic, _) in enumerate(EXPECTED_ITEMS):
            if i >= len(agenda_paras):
                print(f"FAIL: Component 1 — missing paragraph {i+1}")
                continue
            para = agenda_paras[i]
            nonempty_runs = [r for r in para.runs if r.text.strip()]
            if len(nonempty_runs) < 2:
                print(f"FAIL: Component 1 — paragraph {i+1} has fewer than 2 non-empty runs")
                continue
            num_run = nonempty_runs[0]
            topic_run = nonempty_runs[-1]  # last non-empty run is the topic
            if num_run.text.strip() == exp_num and exp_topic.lower() in topic_run.text.strip().lower():
                matched_items += 1
            else:
                print(f"FAIL: Component 1 — item {i+1}: expected '{exp_num}' + '{exp_topic}', "
                      f"found '{num_run.text.strip()}' + '{topic_run.text.strip()}'")

        if matched_items == 6:
            print(f"PASS: Component 1 — all 6 agenda items have correct number and topic text (0.25 pts)")
            total_score += 0.25
        elif matched_items > 0:
            partial = round(0.25 * matched_items / 6, 3)
            print(f"PARTIAL: Component 1 — {matched_items}/6 items correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — no items matched")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Number runs are 28pt (355600 EMU) and bold (0.25 points)
    try:
        comp2_matches = 0
        for i, (exp_num, _, _) in enumerate(EXPECTED_ITEMS):
            if i >= len(agenda_paras):
                continue
            para = agenda_paras[i]
            nonempty_runs = [r for r in para.runs if r.text.strip()]
            if not nonempty_runs:
                continue
            num_run = nonempty_runs[0]
            size_ok = num_run.font.size is not None and num_run.font.size == Pt(28)
            bold_ok = num_run.font.bold is True
            if size_ok and bold_ok:
                comp2_matches += 1
            else:
                print(f"FAIL: Component 2 — item {i+1} number run: size={num_run.font.size} "
                      f"(expected {Pt(28)}), bold={num_run.font.bold} (expected True)")

        if comp2_matches == 6:
            print(f"PASS: Component 2 — all 6 number runs are 28pt bold (0.25 pts)")
            total_score += 0.25
        elif comp2_matches > 0:
            partial = round(0.25 * comp2_matches / 6, 3)
            print(f"PARTIAL: Component 2 — {comp2_matches}/6 number runs correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — no number runs have correct size/bold")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Number colors alternate teal (#008080) / coral (#FF6F61) (0.25 points)
    try:
        comp3_matches = 0
        for i, (_, _, exp_color) in enumerate(EXPECTED_ITEMS):
            if i >= len(agenda_paras):
                continue
            para = agenda_paras[i]
            nonempty_runs = [r for r in para.runs if r.text.strip()]
            if not nonempty_runs:
                continue
            num_run = nonempty_runs[0]
            actual_color = get_run_color_hex(num_run)
            if actual_color and actual_color.upper() == exp_color.upper():
                comp3_matches += 1
            else:
                print(f"FAIL: Component 3 — item {i+1} number color: {actual_color} (expected {exp_color})")

        if comp3_matches == 6:
            print(f"PASS: Component 3 — all 6 number colors alternate correctly (0.25 pts)")
            total_score += 0.25
        elif comp3_matches > 0:
            partial = round(0.25 * comp3_matches / 6, 3)
            print(f"PARTIAL: Component 3 — {comp3_matches}/6 colors correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — no number colors match expected pattern")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Topic runs are 18pt (228600 EMU) and not bold (0.25 points)
    try:
        comp4_matches = 0
        for i, (_, exp_topic, _) in enumerate(EXPECTED_ITEMS):
            if i >= len(agenda_paras):
                continue
            para = agenda_paras[i]
            nonempty_runs = [r for r in para.runs if r.text.strip()]
            if len(nonempty_runs) < 2:
                continue
            topic_run = nonempty_runs[-1]
            size_ok = topic_run.font.size is not None and topic_run.font.size == Pt(18)
            # bold None or False both count as "not bold"
            bold_val = topic_run.font.bold
            bold_ok = bold_val is None or bold_val is False
            if size_ok and bold_ok:
                comp4_matches += 1
            else:
                print(f"FAIL: Component 4 — item {i+1} topic run: size={topic_run.font.size} "
                      f"(expected {Pt(18)}), bold={bold_val} (expected None/False)")

        if comp4_matches == 6:
            print(f"PASS: Component 4 — all 6 topic runs are 18pt regular (0.25 pts)")
            total_score += 0.25
        elif comp4_matches > 0:
            partial = round(0.25 * comp4_matches / 6, 3)
            print(f"PARTIAL: Component 4 — {comp4_matches}/6 topic runs correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — no topic runs have correct size/bold")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
