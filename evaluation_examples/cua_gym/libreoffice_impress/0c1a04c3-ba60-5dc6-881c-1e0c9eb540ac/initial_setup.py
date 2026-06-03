"""
Initial Setup: Create a 10-slide Quarterly Review presentation with empty Agenda body on slide 2.
Task ID: impress_rp_032
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_032'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a styled textbox to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---------- Slide 1: Title Slide ----------
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Quarterly Business Review"
    slide1.placeholders[1].text = "FY2025 Q1 Performance & Strategic Outlook"

    # ---------- Slide 2: Agenda (EMPTY BODY — task target) ----------
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    # Add title manually
    add_textbox(slide2, Inches(0.7), Inches(0.4), Inches(11.9), Inches(1.0),
                "Agenda", font_size=36, bold=True,
                color=RGBColor(0x2E, 0x4A, 0x62), alignment=PP_ALIGN.LEFT)
    # Body area intentionally left empty — no list, no items

    # ---------- Slide 3: Executive Summary ----------
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide3, Inches(0.7), Inches(0.4), Inches(11.9), Inches(1.0),
                "Executive Summary", font_size=32, bold=True,
                color=RGBColor(0x2E, 0x4A, 0x62))
    add_textbox(slide3, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.0),
                "Q1 2025 delivered strong results across all business units. "
                "Revenue grew 14% year-over-year to $127.3M, driven by enterprise "
                "adoption and new product launches. Operating margin improved by "
                "2.1 percentage points to 23.8%. Customer retention remained above "
                "95% for the sixth consecutive quarter.", font_size=18)

    # ---------- Slide 4: Revenue Overview ----------
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide4, Inches(0.7), Inches(0.4), Inches(11.9), Inches(1.0),
                "Q1 Revenue Overview", font_size=32, bold=True,
                color=RGBColor(0x2E, 0x4A, 0x62))
    # Revenue data table-like content
    txBox = slide4.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    revenue_items = [
        ("Enterprise Solutions", "$58.2M", "+18%"),
        ("SMB Products", "$34.7M", "+11%"),
        ("Professional Services", "$22.1M", "+9%"),
        ("Partner Channel", "$12.3M", "+22%"),
    ]
    for i, (segment, amount, growth) in enumerate(revenue_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{segment}:  {amount}  ({growth} YoY)"
        p.space_after = Pt(12)
        run = p.runs[0]
        run.font.size = Pt(18)

    # ---------- Slide 5: Market Analysis ----------
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide5, Inches(0.7), Inches(0.4), Inches(11.9), Inches(1.0),
                "Market Analysis", font_size=32, bold=True,
                color=RGBColor(0x2E, 0x4A, 0x62))
    add_textbox(slide5, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.0),
                "The addressable market expanded to $4.2B in Q1, reflecting a "
                "12% increase from the prior year. Key growth drivers include "
                "digital transformation initiatives across healthcare and financial "
                "services sectors. Competitive intensity remains moderate with our "
                "market share holding steady at 8.3%.", font_size=18)

    # ---------- Slide 6: Customer Metrics ----------
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide6, Inches(0.7), Inches(0.4), Inches(11.9), Inches(1.0),
                "Customer Metrics", font_size=32, bold=True,
                color=RGBColor(0x2E, 0x4A, 0x62))
    txBox6 = slide6.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.0))
    tf6 = txBox6.text_frame
    tf6.word_wrap = True
    metrics = [
        "Total Active Customers: 2,847 (+312 new in Q1)",
        "Net Promoter Score: 72 (up from 68)",
        "Customer Retention Rate: 95.4%",
        "Average Contract Value: $44,700 (+8% YoY)",
        "Support Ticket Resolution: 94% within SLA",
    ]
    for i, m in enumerate(metrics):
        p = tf6.paragraphs[0] if i == 0 else tf6.add_paragraph()
        p.text = m
        p.space_after = Pt(10)
        run = p.runs[0]
        run.font.size = Pt(18)

    # ---------- Slide 7: Product Roadmap ----------
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide7, Inches(0.7), Inches(0.4), Inches(11.9), Inches(1.0),
                "Product Roadmap", font_size=32, bold=True,
                color=RGBColor(0x2E, 0x4A, 0x62))
    add_textbox(slide7, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.0),
                "Q2 priorities include the launch of Analytics Pro 3.0 with "
                "AI-powered insights, mobile app redesign for iOS and Android, "
                "and enterprise SSO integration. The platform migration to "
                "microservices architecture is 67% complete and on track for "
                "Q3 delivery.", font_size=18)

    # ---------- Slide 8: Team Highlights ----------
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide8, Inches(0.7), Inches(0.4), Inches(11.9), Inches(1.0),
                "Team & Hiring Highlights", font_size=32, bold=True,
                color=RGBColor(0x2E, 0x4A, 0x62))
    add_textbox(slide8, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.0),
                "Headcount grew to 438 employees (+23 in Q1). Key hires include "
                "VP of Engineering Sarah Chen, Director of Data Science Raj Patel, "
                "and Senior Product Manager Ayesha Williams. Employee satisfaction "
                "score improved to 4.3/5.0. Engineering velocity increased 15% "
                "through improved CI/CD processes.", font_size=18)

    # ---------- Slide 9: Q2 Goals ----------
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide9, Inches(0.7), Inches(0.4), Inches(11.9), Inches(1.0),
                "Q2 Strategic Goals", font_size=32, bold=True,
                color=RGBColor(0x2E, 0x4A, 0x62))
    txBox9 = slide9.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.0))
    tf9 = txBox9.text_frame
    tf9.word_wrap = True
    goals = [
        "Achieve $135M quarterly revenue target",
        "Launch Analytics Pro 3.0 by June 15",
        "Expand into APAC market with 3 new partners",
        "Reduce customer churn to below 4%",
        "Complete SOC 2 Type II certification",
    ]
    for i, g in enumerate(goals):
        p = tf9.paragraphs[0] if i == 0 else tf9.add_paragraph()
        p.text = g
        p.space_after = Pt(10)
        run = p.runs[0]
        run.font.size = Pt(18)

    # ---------- Slide 10: Thank You ----------
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide10, Inches(2.0), Inches(2.5), Inches(9.0), Inches(2.0),
                "Thank You", font_size=44, bold=True,
                color=RGBColor(0x2E, 0x4A, 0x62), alignment=PP_ALIGN.CENTER)
    add_textbox(slide10, Inches(2.0), Inches(4.0), Inches(9.0), Inches(1.0),
                "Questions & Open Discussion", font_size=24,
                color=RGBColor(0x66, 0x66, 0x66), alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
