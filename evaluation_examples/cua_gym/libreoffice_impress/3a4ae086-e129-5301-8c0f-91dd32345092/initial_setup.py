"""
Initial Setup: Executive Annual Review presentation with 22 slides and 5 sections.
Task ID: impress_gf4_023
Domain: libreoffice_impress

Creates a 22-slide presentation with realistic business content.
NO hyperlinks, navigation buttons, or progress bars.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_023'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Helper to add a styled text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color palette
    DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
    MEDIUM_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
    LIGHT_BLUE = RGBColor(0x3A, 0x7C, 0xBD)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
    ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
    ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)

    section_colors = [
        RGBColor(0x2C, 0x5F, 0x8A),  # Strategy - blue
        RGBColor(0x27, 0xAE, 0x60),  # Operations - green
        RGBColor(0xE6, 0x7E, 0x22),  # Finance - orange
        RGBColor(0x8E, 0x44, 0xAD),  # People - purple
        RGBColor(0xC0, 0x39, 0x2B),  # Outlook - red
    ]

    section_names = ["Strategy", "Operations", "Finance", "People", "Outlook"]

    # ========== SLIDE 1: Title / Main Menu ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    bg = slide1.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BLUE

    add_textbox(slide1, Inches(1), Inches(0.5), Inches(11), Inches(1.5),
                "Executive Annual Review 2025", font_size=40, bold=True,
                color=WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide1, Inches(1), Inches(2.0), Inches(11), Inches(0.7),
                "Meridian Global Partners  |  Board Presentation  |  December 2025",
                font_size=16, color=RGBColor(0xBB, 0xBB, 0xBB), alignment=PP_ALIGN.CENTER)

    # 5 section buttons (text boxes, no hyperlinks)
    button_y = Inches(3.2)
    button_w = Inches(2.0)
    button_h = Inches(1.2)
    gap = Inches(0.3)
    start_x = Inches(1.0)

    for i, name in enumerate(section_names):
        x = start_x + i * (button_w + gap)
        shape = slide1.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            x, button_y, button_w, button_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = section_colors[i]
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = "Calibri"
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = WHITE

    add_textbox(slide1, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
                "Select a section to navigate directly to the topic",
                font_size=14, color=RGBColor(0x99, 0x99, 0x99), alignment=PP_ALIGN.CENTER)

    # ========== SLIDE 2: Agenda / Overview ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide2, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
                "Agenda & Overview", font_size=32, bold=True, color=DARK_BLUE)

    agenda_items = [
        "1. Strategy  -  Corporate direction, market positioning, competitive landscape",
        "2. Operations  -  Process efficiency, supply chain, technology infrastructure",
        "3. Finance  -  Revenue performance, cost management, investment outlook",
        "4. People  -  Talent acquisition, engagement scores, leadership development",
        "5. Outlook  -  2026 targets, risk mitigation, strategic initiatives",
    ]
    y = Inches(1.8)
    for item in agenda_items:
        add_textbox(slide2, Inches(1.0), y, Inches(10), Inches(0.6),
                    item, font_size=18, color=DARK_GRAY)
        y += Inches(0.7)

    add_textbox(slide2, Inches(1.0), Inches(6.0), Inches(10), Inches(0.5),
                "Total presentation time: approximately 90 minutes including Q&A",
                font_size=14, color=RGBColor(0x77, 0x77, 0x77))

    # ========== SECTION 1: STRATEGY (Slides 3-6) ==========
    # Slide 3: Strategy Overview
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s3, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
                "Strategy: Corporate Direction", font_size=32, bold=True, color=section_colors[0])
    add_textbox(s3, Inches(0.8), Inches(1.8), Inches(10), Inches(4.5),
                "Our three-year strategic plan focuses on geographic expansion into Southeast Asian markets, "
                "digital transformation of core business processes, and development of sustainable product lines. "
                "Market analysis indicates a $4.2B addressable market in key growth segments.\n\n"
                "Key strategic pillars:\n"
                "- Expand APAC presence by 35% by Q4 2027\n"
                "- Launch AI-driven logistics platform\n"
                "- Achieve carbon neutrality across manufacturing\n"
                "- Strengthen B2B channel partnerships",
                font_size=16, color=DARK_GRAY)

    # Slide 4: Market Analysis
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s4, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
                "Market Analysis & Competitive Position", font_size=28, bold=True, color=section_colors[0])
    add_textbox(s4, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.0),
                "Market Share by Region (2025):\n\n"
                "North America:  28.3%  (+2.1 YoY)\n"
                "Europe:  22.7%  (+0.8 YoY)\n"
                "Asia-Pacific:  14.5%  (+4.3 YoY)\n"
                "Latin America:  8.2%  (+1.5 YoY)\n"
                "Middle East:  5.1%  (+0.6 YoY)",
                font_size=16, color=DARK_GRAY)
    add_textbox(s4, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.0),
                "Competitive Landscape:\n\n"
                "Top 3 competitors hold 45% combined share.\n"
                "Nexus Corp entered APAC Q2 2025 - aggressive pricing.\n"
                "Our differentiation: integrated supply chain platform + brand trust.\n\n"
                "Threat assessment: MEDIUM\n"
                "Largest gap: Digital self-service portal (launching Q1 2026)",
                font_size=16, color=DARK_GRAY)

    # Slide 5: Strategic Initiatives
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s5, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
                "Strategic Initiatives Pipeline", font_size=28, bold=True, color=section_colors[0])
    table_shape = s5.shapes.add_table(6, 4, Inches(0.8), Inches(1.6), Inches(11), Inches(4.0))
    table = table_shape.table
    headers = ["Initiative", "Investment", "Timeline", "Expected ROI"]
    data = [
        ["APAC Hub Launch", "$18.5M", "Q2 2026 - Q4 2026", "3.2x by 2028"],
        ["AI Logistics Platform", "$12.3M", "Q1 2026 - Q3 2027", "2.8x by 2029"],
        ["Green Manufacturing", "$8.7M", "Q3 2026 - Q2 2028", "1.9x by 2030"],
        ["Digital Portal", "$5.4M", "Q1 2026 - Q2 2026", "4.1x by 2027"],
        ["Partnership Program", "$3.2M", "Ongoing", "2.5x annually"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = section_colors[0]
    for r, row in enumerate(data, 1):
        for c, val in enumerate(row):
            table.cell(r, c).text = val

    # Slide 6: Strategy Summary
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s6, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
                "Strategy: Key Takeaways", font_size=28, bold=True, color=section_colors[0])
    add_textbox(s6, Inches(0.8), Inches(1.8), Inches(10), Inches(4.5),
                "Summary:\n\n"
                "- Total strategic investment of $48.1M across five core initiatives\n"
                "- Expected weighted average ROI of 2.9x within 3-year horizon\n"
                "- APAC expansion remains top priority with highest growth potential\n"
                "- Digital transformation will reduce operational costs by an estimated 15-20%\n"
                "- Board approval requested for Phase 1 funding ($23.9M) by January 2026",
                font_size=16, color=DARK_GRAY)

    # ========== SECTION 2: OPERATIONS (Slides 7-10) ==========
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s7, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
                "Operations: Performance Dashboard", font_size=32, bold=True, color=section_colors[1])
    add_textbox(s7, Inches(0.8), Inches(1.8), Inches(10), Inches(4.5),
                "Operational KPIs (2025 YTD):\n\n"
                "Order Fulfillment Rate: 97.3% (target: 96.0%)\n"
                "Average Delivery Time: 2.4 days (target: 3.0 days)\n"
                "Warehouse Utilization: 84.7% (target: 85.0%)\n"
                "Defect Rate: 0.42% (target: 0.50%)\n"
                "On-Time Shipping: 95.1% (target: 94.0%)\n\n"
                "Overall operations score: 94.8/100 - exceeding annual targets",
                font_size=16, color=DARK_GRAY)

    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s8, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
                "Supply Chain Resilience", font_size=28, bold=True, color=section_colors[1])
    add_textbox(s8, Inches(0.8), Inches(1.8), Inches(10), Inches(4.5),
                "Supply chain diversification efforts yielded measurable improvements:\n\n"
                "- Reduced single-supplier dependency from 34% to 18%\n"
                "- Added 12 new qualified suppliers across 4 regions\n"
                "- Implemented real-time tracking across 89% of shipments\n"
                "- Buffer stock optimization saved $2.3M in carrying costs\n"
                "- Average lead time reduced from 14.2 to 11.5 days\n\n"
                "Risk exposure index decreased from 7.2 to 4.8 (scale 1-10)",
                font_size=16, color=DARK_GRAY)

    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s9, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
                "Technology Infrastructure Modernization", font_size=28, bold=True, color=section_colors[1])
    add_textbox(s9, Inches(0.8), Inches(1.8), Inches(10), Inches(4.5),
                "Cloud migration progress: 72% complete (target: 80% by Q2 2026)\n\n"
                "Key technology upgrades:\n"
                "- ERP system upgrade to SAP S/4HANA (completed Aug 2025)\n"
                "- Warehouse management automation (Phase 2 in progress)\n"
                "- Predictive maintenance AI deployed at 3 manufacturing sites\n"
                "- Cybersecurity framework upgraded to NIST 2.0 standards\n\n"
                "IT spend as % of revenue: 4.8% (industry avg: 5.2%)",
                font_size=16, color=DARK_GRAY)

    s10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s10, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
                "Operations: Key Takeaways", font_size=28, bold=True, color=section_colors[1])
    add_textbox(s10, Inches(0.8), Inches(1.8), Inches(10), Inches(4.5),
                "Summary:\n\n"
                "- Operational efficiency improved 8.2% year-over-year\n"
                "- Supply chain risk score at historic low of 4.8\n"
                "- Technology modernization on track, 72% cloud migration complete\n"
                "- Projected operational savings of $6.1M in 2026\n"
                "- Automation investments showing 18-month payback period",
                font_size=16, color=DARK_GRAY)

    # ========== SECTION 3: FINANCE (Slides 11-14) ==========
    s11 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s11, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
                "Finance: Revenue & Profitability", font_size=32, bold=True, color=section_colors[2])
    add_textbox(s11, Inches(0.8), Inches(1.8), Inches(10), Inches(4.5),
                "FY 2025 Financial Highlights:\n\n"
                "Total Revenue: $847.3M (+12.4% YoY)\n"
                "Gross Margin: 42.1% (up from 39.8%)\n"
                "EBITDA: $178.6M (21.1% margin)\n"
                "Net Income: $112.4M (+18.7% YoY)\n"
                "Free Cash Flow: $94.2M\n"
                "Earnings Per Share: $4.82 (vs. $4.06 prior year)\n\n"
                "Revenue exceeded guidance by $23M, driven by APAC growth",
                font_size=16, color=DARK_GRAY)

    s12 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s12, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
                "Cost Structure & Efficiency", font_size=28, bold=True, color=section_colors[2])
    table_shape2 = s12.shapes.add_table(7, 3, Inches(0.8), Inches(1.6), Inches(10), Inches(4.5))
    t2 = table_shape2.table
    h2 = ["Cost Category", "FY 2025 ($M)", "% of Revenue"]
    d2 = [
        ["Cost of Goods Sold", "490.2", "57.9%"],
        ["R&D", "67.8", "8.0%"],
        ["Sales & Marketing", "84.7", "10.0%"],
        ["G&A", "42.4", "5.0%"],
        ["Depreciation & Amortization", "33.9", "4.0%"],
        ["Other Operating Expenses", "16.9", "2.0%"],
    ]
    for c, h in enumerate(h2):
        cell = t2.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = section_colors[2]
    for r, row in enumerate(d2, 1):
        for c, val in enumerate(row):
            t2.cell(r, c).text = val

    s13 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s13, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
                "Capital Allocation & Investment", font_size=28, bold=True, color=section_colors[2])
    add_textbox(s13, Inches(0.8), Inches(1.8), Inches(10), Inches(4.5),
                "Capital allocation priorities for 2026:\n\n"
                "1. Organic growth investments: $48.1M (strategic initiatives)\n"
                "2. Share buyback program: $30.0M (continuing)\n"
                "3. Dividend increase: 8% to $1.92/share annually\n"
                "4. M&A reserve: $25.0M (targeting complementary technologies)\n"
                "5. Debt reduction: $15.0M (targeting 1.8x leverage ratio)\n\n"
                "Total capital deployment plan: $118.1M\n"
                "Balance sheet remains strong with $142M cash and equivalents",
                font_size=16, color=DARK_GRAY)

    s14 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s14, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
                "Finance: Key Takeaways", font_size=28, bold=True, color=section_colors[2])
    add_textbox(s14, Inches(0.8), Inches(1.8), Inches(10), Inches(4.5),
                "Summary:\n\n"
                "- Revenue up 12.4% to $847.3M, exceeding guidance\n"
                "- Margin expansion of 230bps driven by operational efficiency\n"
                "- Strong free cash flow of $94.2M supports growth agenda\n"
                "- 2026 guidance: Revenue $920-$950M, EBITDA margin 22-23%\n"
                "- Analyst consensus price target implies 18% upside",
                font_size=16, color=DARK_GRAY)

    # ========== SECTION 4: PEOPLE (Slides 15-18) ==========
    s15 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s15, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
                "People: Talent & Culture", font_size=32, bold=True, color=section_colors[3])
    add_textbox(s15, Inches(0.8), Inches(1.8), Inches(10), Inches(4.5),
                "Workforce Metrics (2025):\n\n"
                "Total Headcount: 4,287 (+8.3% YoY)\n"
                "Employee Engagement Score: 82/100 (industry avg: 74)\n"
                "Voluntary Turnover: 11.2% (industry avg: 15.8%)\n"
                "Internal Mobility Rate: 24.3%\n"
                "Diversity Index: 0.78 (up from 0.71)\n"
                "Average Tenure: 4.7 years\n\n"
                "Named Top 50 Employer by Workplace Excellence Institute",
                font_size=16, color=DARK_GRAY)

    s16 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s16, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
                "Talent Acquisition & Development", font_size=28, bold=True, color=section_colors[3])
    add_textbox(s16, Inches(0.8), Inches(1.8), Inches(10), Inches(4.5),
                "Recruitment:\n"
                "- 628 new hires in 2025 (68% met quality benchmark)\n"
                "- Average time-to-fill: 32 days (down from 41 days)\n"
                "- Campus recruiting pipeline expanded to 18 universities\n\n"
                "Learning & Development:\n"
                "- $3.4M invested in training programs\n"
                "- 92% completion rate for mandatory compliance training\n"
                "- 340 employees enrolled in leadership acceleration program\n"
                "- New mentorship platform launched with 78% participation",
                font_size=16, color=DARK_GRAY)

    s17 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s17, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
                "Compensation & Benefits Review", font_size=28, bold=True, color=section_colors[3])
    add_textbox(s17, Inches(0.8), Inches(1.8), Inches(10), Inches(4.5),
                "Compensation benchmarking completed Q3 2025:\n\n"
                "- Base pay competitiveness ratio: 1.03 (at market)\n"
                "- Total compensation ratio: 1.08 (above market)\n"
                "- Merit increase budget for 2026: 4.2%\n"
                "- Executive incentive plan alignment: 75% long-term/25% short-term\n"
                "- Benefits cost per employee: $14,200 annually\n\n"
                "New benefits introduced: expanded parental leave (16 weeks), "
                "mental health stipend ($1,500/year), flexible work arrangement policy",
                font_size=16, color=DARK_GRAY)

    s18 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s18, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
                "People: Key Takeaways", font_size=28, bold=True, color=section_colors[3])
    add_textbox(s18, Inches(0.8), Inches(1.8), Inches(10), Inches(4.5),
                "Summary:\n\n"
                "- Engagement score at 82/100, highest in company history\n"
                "- Turnover well below industry average at 11.2%\n"
                "- Diversity index improved to 0.78, on track for 0.82 target\n"
                "- Leadership pipeline strengthened with 340 participants\n"
                "- 2026 focus: AI skills development and succession planning",
                font_size=16, color=DARK_GRAY)

    # ========== SECTION 5: OUTLOOK (Slides 19-22) ==========
    s19 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s19, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
                "Outlook: 2026 Targets", font_size=32, bold=True, color=section_colors[4])
    add_textbox(s19, Inches(0.8), Inches(1.8), Inches(10), Inches(4.5),
                "Key Performance Targets for 2026:\n\n"
                "Revenue: $920M - $950M (+8.6% to +12.1%)\n"
                "EBITDA Margin: 22% - 23%\n"
                "Net Income: $125M - $135M\n"
                "Free Cash Flow: $100M+\n"
                "Market Share (Global): 20%+\n"
                "Customer Satisfaction: 90+ NPS\n"
                "Employee Engagement: 84+\n\n"
                "These targets assume stable macroeconomic conditions and no major disruptions",
                font_size=16, color=DARK_GRAY)

    s20 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s20, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
                "Risk Assessment & Mitigation", font_size=28, bold=True, color=section_colors[4])
    table_shape3 = s20.shapes.add_table(6, 4, Inches(0.8), Inches(1.6), Inches(11), Inches(4.2))
    t3 = table_shape3.table
    h3 = ["Risk Factor", "Probability", "Impact", "Mitigation"]
    d3 = [
        ["Economic downturn", "Medium", "High", "Diversified revenue, cost flex plan"],
        ["Supply chain disruption", "Medium", "Medium", "Multi-source strategy, buffer stock"],
        ["Cybersecurity breach", "Low", "Critical", "NIST 2.0, insurance, incident plan"],
        ["Talent market tightening", "High", "Medium", "Retention programs, employer brand"],
        ["Regulatory changes", "Medium", "Medium", "Compliance team, government affairs"],
    ]
    for c, h in enumerate(h3):
        cell = t3.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(13)
            run.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = section_colors[4]
    for r, row in enumerate(d3, 1):
        for c, val in enumerate(row):
            t3.cell(r, c).text = val

    s21 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s21, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
                "Strategic Roadmap 2026-2028", font_size=28, bold=True, color=section_colors[4])
    add_textbox(s21, Inches(0.8), Inches(1.8), Inches(10), Inches(4.5),
                "Phase 1 (H1 2026): Foundation\n"
                "- Launch APAC regional hub in Singapore\n"
                "- Deploy AI logistics platform (pilot with 3 key accounts)\n"
                "- Complete cloud migration to 80%\n\n"
                "Phase 2 (H2 2026): Acceleration\n"
                "- Scale digital portal to all B2B customers\n"
                "- Begin green manufacturing retrofit at 2 plants\n"
                "- Expand partner network by 25 new agreements\n\n"
                "Phase 3 (2027-2028): Optimization & Scale\n"
                "- Full APAC market penetration target: 20% share\n"
                "- AI-driven operations across all facilities\n"
                "- Carbon neutrality milestone for manufacturing",
                font_size=16, color=DARK_GRAY)

    s22 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s22, Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
                "Outlook: Closing & Next Steps", font_size=28, bold=True, color=section_colors[4])
    add_textbox(s22, Inches(0.8), Inches(1.8), Inches(10), Inches(4.5),
                "Meridian Global Partners is positioned for sustained growth.\n\n"
                "Immediate next steps requiring board action:\n"
                "1. Approve Phase 1 strategic investment ($23.9M)\n"
                "2. Ratify 2026 compensation framework and merit budget\n"
                "3. Authorize share buyback continuation ($30M)\n"
                "4. Confirm Q1 2026 board meeting schedule\n\n"
                "Thank you for your time and continued guidance.\n"
                "Questions and discussion to follow.",
                font_size=16, color=DARK_GRAY)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
