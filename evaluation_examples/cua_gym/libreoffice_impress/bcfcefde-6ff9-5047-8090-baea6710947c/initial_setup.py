"""
Initial Setup: Corporate annual report cover slide - blank starting state
Task ID: impress_gf2_036
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_036'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard slide size (default 10x7.5 inches) - leave as is
    # Add one blank slide with white background (no shapes/text)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only layout
    # Remove the title placeholder so it's truly empty
    for shape in list(slide.placeholders):
        sp = shape.element
        sp.getparent().remove(sp)

    # Set explicit white background
    fill = slide.background.fill
    fill.solid()
    from pptx.dml.color import RGBColor
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
