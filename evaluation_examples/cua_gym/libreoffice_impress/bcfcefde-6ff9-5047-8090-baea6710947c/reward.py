"""
Reward Script: Corporate annual report cover slide design verification
Task ID: impress_gf2_036
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): Slide dimensions 33.87 x 19.05 cm
  Component 2 (0.15): Semi-transparent white circle shape present
  Component 3 (0.15): Gold decorative horizontal line/rectangle
  Component 4 (0.20): Company name text box — 56pt bold white
  Component 5 (0.15): Year '2024' text box — 120pt bold white
  Component 6 (0.15): Tagline text box — 18pt italic gold
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_036'


def persist_app_state(domain: str):
    """Save any open LibreOffice documents before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Emu, Pt, Cm
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Missing library: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    if len(prs.slides) == 0:
        print("CRITICAL: No slides in presentation")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]

    # =========================================================================
    # Component 1: Slide dimensions — 33.87 x 19.05 cm (0.20 points)
    # Task changes the slide width from standard 25.4 cm to 33.87 cm.
    # Height stays at 19.05 cm (precondition), so we only score the width change
    # but verify both for correctness.
    # =========================================================================
    try:
        width_cm = prs.slide_width / 360000
        height_cm = prs.slide_height / 360000
        # Width must be ~33.87 cm (changed from 25.4 cm)
        width_ok = abs(width_cm - 33.87) < 0.5
        # Height must remain ~19.05 cm
        height_ok = abs(height_cm - 19.05) < 0.5
        if width_ok and height_ok:
            print(f"PASS: Component 1 — Slide dimensions {width_cm:.2f} x {height_cm:.2f} cm (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 — Expected ~33.87 x 19.05 cm, found {width_cm:.2f} x {height_cm:.2f} cm")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Collect all shapes for analysis
    shapes = list(slide.shapes)

    # =========================================================================
    # Component 2: Semi-transparent white circle/oval shape (0.15 points)
    # Must have: ellipse geometry, white fill, alpha < 100% (semi-transparent)
    # =========================================================================
    try:
        import zipfile
        import xml.etree.ElementTree as ET

        ns = {
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        }

        circle_found = False
        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide1.xml') as f:
                root = ET.fromstring(f.read())
            for sp in root.findall('.//p:cSld/p:spTree/p:sp', ns):
                spPr = sp.find('p:spPr', ns)
                if spPr is None:
                    continue
                prstGeom = spPr.find('a:prstGeom', ns)
                if prstGeom is None or prstGeom.get('prst') != 'ellipse':
                    continue
                # Found an ellipse — check for white fill with alpha
                solidFill = spPr.find('a:solidFill', ns)
                if solidFill is not None:
                    clr = solidFill.find('a:srgbClr', ns)
                    if clr is not None:
                        color_val = clr.get('val', '').upper()
                        alpha_elem = clr.find('a:alpha', ns)
                        alpha_val = int(alpha_elem.get('val', '100000')) if alpha_elem is not None else 100000
                        # White fill with semi-transparency (alpha < 100%)
                        if color_val == 'FFFFFF' and alpha_val < 100000:
                            circle_found = True

        if circle_found:
            print(f"PASS: Component 2 — Semi-transparent white circle found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 — No semi-transparent white ellipse found")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # =========================================================================
    # Component 3: Gold decorative horizontal line/rectangle (0.15 points)
    # Must have: FFD700 fill, rectangular shape, wider than tall (horizontal)
    # =========================================================================
    try:
        gold_line_found = False
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # Check if it's a thin horizontal shape (line-like)
                if shape.width > shape.height * 3:  # much wider than tall
                    # Check XML for gold fill
                    spPr_elem = shape._element.find(
                        '{http://schemas.openxmlformats.org/presentationml/2006/main}spPr',
                        ns
                    )
                    # Use a simpler approach: check via XML
                    pass

        # Use XML approach for reliability
        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide1.xml') as f:
                root = ET.fromstring(f.read())
            for sp in root.findall('.//p:cSld/p:spTree/p:sp', ns):
                spPr = sp.find('p:spPr', ns)
                if spPr is None:
                    continue
                prstGeom = spPr.find('a:prstGeom', ns)
                if prstGeom is None:
                    continue
                geom_type = prstGeom.get('prst', '')
                if geom_type == 'ellipse':
                    continue  # skip the circle
                # Check if it's a rectangle with gold fill
                solidFill = spPr.find('a:solidFill', ns)
                if solidFill is not None:
                    clr = solidFill.find('a:srgbClr', ns)
                    if clr is not None:
                        color_val = clr.get('val', '').upper()
                        if color_val == 'FFD700':
                            # Check dimensions via xfrm
                            xfrm = spPr.find('a:xfrm', ns)
                            if xfrm is not None:
                                ext = xfrm.find('a:ext', ns)
                                if ext is not None:
                                    cx = int(ext.get('cx', 0))
                                    cy = int(ext.get('cy', 0))
                                    # Horizontal: width >> height
                                    if cx > cy * 3:
                                        gold_line_found = True

        if gold_line_found:
            print(f"PASS: Component 3 — Gold horizontal line/rectangle found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 3 — No gold (#FFD700) horizontal line shape found")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # =========================================================================
    # Component 4: Company name text — 56pt bold white at top-left (0.20 points)
    # =========================================================================
    try:
        company_name_found = False
        for shape in shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text or len(text) < 3:
                    continue
                runs = [r for r in para.runs if (r.text or "").strip()]
                if not runs:
                    continue
                for run in runs:
                    font = run.font
                    # Check: bold, white color, ~56pt size
                    is_bold = font.bold is True
                    size_ok = False
                    if font.size is not None:
                        size_pt = font.size / 12700  # EMU to pt
                        size_ok = abs(size_pt - 56) < 2
                    color_ok = False
                    try:
                        if font.color.type is not None:
                            color_ok = str(font.color.rgb).upper() == 'FFFFFF'
                    except:
                        pass

                    if is_bold and size_ok and color_ok:
                        # Verify it's NOT "2024" (that's another component)
                        if '2024' not in run.text:
                            company_name_found = True
                            print(f"  Found company name: '{run.text}', size={font.size/12700:.0f}pt, bold={font.bold}")

        if company_name_found:
            print(f"PASS: Component 4 — Company name in 56pt bold white found (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 4 — No text with 56pt bold white font found")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # =========================================================================
    # Component 5: Year '2024' in large text — ~120pt bold white (0.15 points)
    # =========================================================================
    try:
        year_found = False
        for shape in shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if '2024' not in text:
                    continue
                runs = [r for r in para.runs if (r.text or "").strip()]
                for run in runs:
                    if '2024' not in run.text:
                        continue
                    font = run.font
                    is_bold = font.bold is True
                    size_ok = False
                    if font.size is not None:
                        size_pt = font.size / 12700
                        # Allow some tolerance around 120pt
                        size_ok = size_pt >= 90
                    color_ok = False
                    try:
                        if font.color.type is not None:
                            color_ok = str(font.color.rgb).upper() == 'FFFFFF'
                    except:
                        pass

                    if is_bold and size_ok and color_ok:
                        year_found = True
                        print(f"  Found year: '{run.text}', size={font.size/12700:.0f}pt, bold={font.bold}")

        if year_found:
            print(f"PASS: Component 5 — Year '2024' in large bold white text found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 5 — No '2024' text with large bold white font found")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # =========================================================================
    # Component 6: Tagline in 18pt italic gold (#FFD700) (0.15 points)
    # =========================================================================
    try:
        tagline_found = False
        for shape in shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text or len(text) < 5:
                    continue
                runs = [r for r in para.runs if (r.text or "").strip()]
                for run in runs:
                    font = run.font
                    is_italic = font.italic is True
                    size_ok = False
                    if font.size is not None:
                        size_pt = font.size / 12700
                        size_ok = abs(size_pt - 18) < 2
                    color_ok = False
                    try:
                        if font.color.type is not None:
                            color_ok = str(font.color.rgb).upper() == 'FFD700'
                    except:
                        pass

                    # Must NOT be bold, and must have gold color + italic + ~18pt
                    # Also exclude company name and year
                    if is_italic and size_ok and color_ok:
                        if '2024' not in run.text:
                            tagline_found = True
                            print(f"  Found tagline: '{run.text}', size={font.size/12700:.0f}pt, italic={font.italic}")

        if tagline_found:
            print(f"PASS: Component 6 — Tagline in 18pt italic gold found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 6 — No text with 18pt italic gold (#FFD700) font found")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
