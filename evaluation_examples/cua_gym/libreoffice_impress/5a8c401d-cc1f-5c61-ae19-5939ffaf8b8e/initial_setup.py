"""
Initial Setup: Create a 12-slide Annual Review 2025 presentation
Task ID: impress_ps_041
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_041'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.paragraphs[0].text = point
        else:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(prs,
        "Annual Review 2025",
        "Acme Corporation | Presented by Sarah Chen, CEO"
    )

    # Slide 2: Agenda
    add_content_slide(prs, "Agenda", [
        "Company Overview & Mission Update",
        "Financial Performance Summary",
        "Product & Engineering Milestones",
        "Revenue Breakdown by Division",
        "Market Position & Competitive Analysis",
        "Team Growth & Organizational Changes",
        "Key Initiatives for 2026",
    ])

    # Slide 3: Company Overview
    add_content_slide(prs, "Company Overview", [
        "Founded in 2018 with a mission to democratize AI-powered analytics",
        "Expanded to 12 global offices across North America, Europe, and Asia",
        "Grew customer base from 2,400 to 4,150 enterprise accounts (+73%)",
        "Net Promoter Score improved from 62 to 78",
        "Named in Forbes Top 50 AI Companies for the third consecutive year",
    ])

    # Slide 4: Financial Highlights
    add_content_slide(prs, "Financial Highlights", [
        "Total Revenue: $287.4M (up 42% YoY)",
        "Gross Margin: 74.2% (up from 71.8%)",
        "Operating Income: $38.6M",
        "Free Cash Flow: $52.1M",
        "ARR reached $312M by December 2025",
    ])

    # Slide 5: Revenue Breakdown
    add_content_slide(prs, "Revenue Breakdown by Division", [
        "Enterprise Solutions: $142.3M (49.5%)",
        "SMB Platform: $68.9M (24.0%)",
        "Professional Services: $41.2M (14.3%)",
        "Government & Public Sector: $35.0M (12.2%)",
        "Average contract value increased 28% to $69,200",
    ])

    # Slide 6: Market Position
    add_content_slide(prs, "Market Position & Strategy", [
        "Ranked #2 in Gartner Magic Quadrant for Analytics Platforms",
        "Market share grew from 8.3% to 11.7%",
        "Strategic partnerships with AWS, Azure, and Google Cloud",
        "Launched Acme AI Copilot with 1,200+ beta users",
        "Patent portfolio expanded to 47 granted patents",
    ])

    # Slide 7: Product Milestones
    add_content_slide(prs, "Product & Engineering Milestones", [
        "Released Acme Platform v5.0 with real-time collaboration",
        "Reduced average query response time by 63%",
        "Shipped 142 feature updates across quarterly releases",
        "Mobile app downloads surpassed 500,000",
        "99.97% platform uptime maintained throughout 2025",
    ])

    # Slide 8: Team Overview
    add_content_slide(prs, "Team Growth & Culture", [
        "Total headcount: 1,247 employees (up from 892)",
        "Engineering team expanded to 485 members",
        "Employee retention rate: 91.4%",
        "Diversity: 44% women in leadership roles",
        "Launched internal mentorship program with 320 participants",
    ])

    # Slide 9: Department Performance
    add_content_slide(prs, "Department Performance Ratings", [
        "Engineering: 4.6/5.0 - Exceeded delivery targets by 18%",
        "Sales: 4.4/5.0 - Closed 287 new enterprise accounts",
        "Marketing: 4.3/5.0 - Generated 12,400 qualified leads",
        "Customer Success: 4.7/5.0 - Churn reduced to 3.2%",
        "HR & Operations: 4.2/5.0 - Streamlined onboarding to 5 days",
    ])

    # Slide 10: Leadership Changes
    add_content_slide(prs, "Leadership & Organizational Updates", [
        "Appointed Marcus Johnson as CTO (previously VP Engineering)",
        "Created Chief Data Officer role - filled by Dr. Priya Patel",
        "Restructured into 4 business units for faster decision-making",
        "Promoted 23 employees to senior leadership positions",
        "Established AI Ethics Board with 5 external advisors",
    ])

    # Slide 11: Employee Engagement
    add_content_slide(prs, "Employee Engagement & Wellbeing", [
        "Annual engagement score: 8.4/10 (industry avg: 7.1)",
        "Introduced 4-day work week pilot for Q3 (92% approval)",
        "Expanded mental health benefits to include family coverage",
        "Launched Learning & Development budget of $3,500/employee",
        "156 employees completed leadership development program",
    ])

    # Slide 12: Looking Ahead
    add_content_slide(prs, "2026 Strategic Priorities", [
        "Target $420M revenue with 46% YoY growth",
        "Launch Acme Platform v6.0 with generative AI features",
        "Expand into Japan and Southeast Asia markets",
        "Achieve SOC 2 Type II and ISO 27001 certifications",
        "Grow headcount to 1,800 with focus on AI research",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
