"""
Initial Setup: Insert image onto slide 2 of presentation
Task ID: impress_tm_046
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_046'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
PHOTO_DIR = f'{WORKDIR}/photos'
PHOTO_PATH = f'{PHOTO_DIR}/team_photo.jpg'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_team_photo():
    """Create a realistic-looking 1920x1080 team photo image."""
    os.makedirs(PHOTO_DIR, exist_ok=True)
    # Create a gradient image that looks like a team photo backdrop
    img = Image.new('RGB', (1920, 1080))
    pixels = img.load()
    for y in range(1080):
        for x in range(1920):
            # Warm gradient simulating an office/team setting
            r = min(255, int(60 + (x / 1920) * 120 + (y / 1080) * 40))
            g = min(255, int(80 + (x / 1920) * 80 + (y / 1080) * 30))
            b = min(255, int(120 + (x / 1920) * 60 + (y / 1080) * 50))
            pixels[x, y] = (r, g, b)
    img.save(PHOTO_PATH, 'JPEG', quality=85)
    print(f'Team photo created: {PHOTO_PATH}')


def create_initial():
    prs = Presentation()
    # Standard widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
    slide1.shapes.title.text = "Meridian Technologies"
    slide1.placeholders[1].text = "Annual Company Overview 2025"

    # --- Slide 2: "Our Team" - Title Only (NO image, NO content placeholder) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Our Team"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x6E)

    # --- Slide 3: Our Projects ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox3_title = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
    tf3t = txBox3_title.text_frame
    p3t = tf3t.paragraphs[0]
    p3t.text = "Our Projects"
    p3t.alignment = PP_ALIGN.LEFT
    r3t = p3t.runs[0]
    r3t.font.name = "Calibri"
    r3t.font.size = Pt(36)
    r3t.font.bold = True
    r3t.font.color.rgb = RGBColor(0x2E, 0x4A, 0x6E)

    # Project content
    txBox3_body = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(5))
    tf3b = txBox3_body.text_frame
    tf3b.word_wrap = True
    projects = [
        ("Project Aurora", "Cloud migration platform - Phase 3 deployment underway"),
        ("Project Horizon", "Next-gen analytics dashboard - Beta testing with key clients"),
        ("Project Catalyst", "AI-driven process automation - Reducing ops costs by 35%"),
        ("Project Nexus", "Cross-platform integration suite - Launch scheduled Q3 2025"),
    ]
    for i, (name, desc) in enumerate(projects):
        if i > 0:
            p = tf3b.add_paragraph()
        else:
            p = tf3b.paragraphs[0]
        p.text = f"{name}"
        run = p.runs[0]
        run.font.name = "Calibri"
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        p2 = tf3b.add_paragraph()
        p2.text = desc
        r2 = p2.runs[0]
        r2.font.name = "Calibri"
        r2.font.size = Pt(16)
        r2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

        # Add spacing between projects
        if i < len(projects) - 1:
            spacer = tf3b.add_paragraph()
            spacer.text = ""

    # --- Slide 4: Contact Us ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox4_title = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
    tf4t = txBox4_title.text_frame
    p4t = tf4t.paragraphs[0]
    p4t.text = "Contact Us"
    p4t.alignment = PP_ALIGN.LEFT
    r4t = p4t.runs[0]
    r4t.font.name = "Calibri"
    r4t.font.size = Pt(36)
    r4t.font.bold = True
    r4t.font.color.rgb = RGBColor(0x2E, 0x4A, 0x6E)

    txBox4_body = slide4.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11), Inches(4))
    tf4b = txBox4_body.text_frame
    tf4b.word_wrap = True
    contacts = [
        "Email: info@meridiantech.com",
        "Phone: +1 (555) 234-8900",
        "Address: 450 Innovation Drive, Suite 1200, San Francisco, CA 94105",
        "Website: www.meridiantech.com",
    ]
    for i, line in enumerate(contacts):
        if i > 0:
            p = tf4b.add_paragraph()
        else:
            p = tf4b.paragraphs[0]
        p.text = line
        run = p.runs[0]
        run.font.name = "Calibri"
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_team_photo()
create_initial()
