"""
Reward Script: Insert team_photo.jpg onto slide 2 of team_intro.pptx
Task ID: impress_tm_046
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Slide 2 contains a PICTURE shape
  Component 2 (0.3): The image blob matches /home/user/photos/team_photo.jpg
  Component 3 (0.3): Image has non-zero dimensions and is within slide bounds
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_046'


def persist_app_state(domain: str):
    """Try to save any unsaved LibreOffice edits via Ctrl+S."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            import time
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 2")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]  # 0-indexed, slide 2

    # Find all PICTURE shapes on slide 2
    picture_shapes = []
    for shape in slide2.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture_shapes.append(shape)

    # Component 1: Slide 2 contains at least one PICTURE shape (0.4 points)
    # This checks the core task: an image was inserted onto slide 2.
    # Initial env has NO pictures on slide 2; golden env has one.
    try:
        if len(picture_shapes) > 0:
            print(f"PASS: Component 1 — Slide 2 has {len(picture_shapes)} picture(s) (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 — Slide 2 has no picture shapes")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: The inserted image matches team_photo.jpg (0.3 points)
    # Verifies the correct image file was used, not an arbitrary image.
    try:
        source_path = os.path.join(WORKDIR, 'photos', 'team_photo.jpg')
        if not os.path.exists(source_path):
            print(f"FAIL: Component 2 — Source image not found at {source_path}")
        elif len(picture_shapes) > 0:
            with open(source_path, 'rb') as f:
                source_blob = f.read()
            matching_blobs = [pic for pic in picture_shapes if pic.image.blob == source_blob]
            if len(matching_blobs) > 0:
                print(f"PASS: Component 2 — Image blob matches team_photo.jpg ({len(source_blob)} bytes) (0.3 pts)")
                total_score += 0.3
            else:
                # Fallback: check size match (image may have been re-encoded)
                matching_sizes = [pic for pic in picture_shapes if len(pic.image.blob) == len(source_blob)]
                if len(matching_sizes) > 0:
                    print(f"PASS: Component 2 — Image size matches team_photo.jpg (0.3 pts)")
                    total_score += 0.3
                else:
                    print(f"FAIL: Component 2 — No picture on slide 2 matches team_photo.jpg blob")
        else:
            print(f"FAIL: Component 2 — No pictures on slide 2 to compare")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Image has valid dimensions and position (0.3 points)
    # Verifies the image is actually visible (non-zero size, within slide area).
    try:
        if len(picture_shapes) > 0:
            pic = picture_shapes[0]
            slide_w = prs.slide_width
            slide_h = prs.slide_height
            has_size = pic.width > 0 and pic.height > 0
            # Check image is at least partially within slide bounds
            within_bounds = (pic.left < slide_w and pic.top < slide_h and
                             pic.left + pic.width > 0 and pic.top + pic.height > 0)
            if has_size and within_bounds:
                print(f"PASS: Component 3 — Image dimensions valid: "
                      f"w={pic.width}, h={pic.height}, left={pic.left}, top={pic.top} (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 3 — Image has invalid dimensions or out of bounds: "
                      f"w={pic.width}, h={pic.height}, left={pic.left}, top={pic.top}")
        else:
            print(f"FAIL: Component 3 — No pictures on slide 2 to check dimensions")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
