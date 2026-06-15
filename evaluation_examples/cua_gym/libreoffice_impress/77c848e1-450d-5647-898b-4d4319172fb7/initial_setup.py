"""
Initial Setup: Create complex_deck.pptx with 8 slides.
Slide 8 has 6 objects added in non-sequential order so tab/accessibility
order is wrong (footer added before content boxes, etc.).
Task ID: impress_gf5_025
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_025'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=None):
    """Helper to add a text box with styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    sl1 = prs.slides.add_slide(prs.slide_layouts[0])
    sl1.shapes.title.text = "Accessibility Compliance Report"
    sl1.placeholders[1].text = "Q1 2026 — Prepared by Digital Accessibility Team"

    # ---- Slide 2: Agenda ----
    sl2 = prs.slides.add_slide(prs.slide_layouts[1])
    sl2.shapes.title.text = "Meeting Agenda"
    body2 = sl2.placeholders[1].text_frame
    body2.text = "1. Current compliance status across products"
    body2.add_paragraph().text = "2. WCAG 2.2 audit results"
    body2.add_paragraph().text = "3. Remediation timeline and milestones"
    body2.add_paragraph().text = "4. Screen reader testing outcomes"
    body2.add_paragraph().text = "5. Training and certification updates"

    # ---- Slide 3: Compliance Overview ----
    sl3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_textbox(sl3, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                "WCAG 2.2 Compliance Overview", font_size=28, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x6B))
    add_textbox(sl3, Inches(0.5), Inches(1.5), Inches(5.5), Inches(5),
                "Product A — Web Portal\n"
                "  Level AA: 94% compliant\n"
                "  Critical issues: 3 remaining\n"
                "  Target date: April 30, 2026\n\n"
                "Product B — Mobile App\n"
                "  Level AA: 87% compliant\n"
                "  Critical issues: 8 remaining\n"
                "  Target date: June 15, 2026",
                font_size=14)
    add_textbox(sl3, Inches(6.5), Inches(1.5), Inches(6), Inches(5),
                "Product C — Internal Dashboard\n"
                "  Level AA: 72% compliant\n"
                "  Critical issues: 14 remaining\n"
                "  Target date: August 1, 2026\n\n"
                "Product D — Customer Portal\n"
                "  Level AA: 91% compliant\n"
                "  Critical issues: 5 remaining\n"
                "  Target date: May 15, 2026",
                font_size=14)

    # ---- Slide 4: Audit Results ----
    sl4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(sl4, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                "WCAG 2.2 Audit Results — March 2026", font_size=28, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x6B))
    table_shape = sl4.shapes.add_table(6, 4, Inches(0.5), Inches(1.5),
                                        Inches(12), Inches(4.5))
    tbl = table_shape.table
    headers = ["Criterion", "Level", "Status", "Notes"]
    data = [
        ["1.1.1 Non-text Content", "A", "Pass", "All images have alt text"],
        ["1.3.1 Info and Relationships", "A", "Fail", "Missing form labels on checkout"],
        ["2.1.1 Keyboard", "A", "Pass", "All interactive elements reachable"],
        ["2.4.7 Focus Visible", "AA", "Partial", "Custom components need focus rings"],
        ["4.1.2 Name, Role, Value", "A", "Fail", "ARIA labels missing on 12 widgets"],
    ]
    for c, h in enumerate(headers):
        tbl.cell(0, c).text = h
    for r, row in enumerate(data, 1):
        for c, val in enumerate(row):
            tbl.cell(r, c).text = val

    # ---- Slide 5: Screen Reader Testing ----
    sl5 = prs.slides.add_slide(prs.slide_layouts[1])
    sl5.shapes.title.text = "Screen Reader Testing Results"
    body5 = sl5.placeholders[1].text_frame
    body5.text = "NVDA (Windows): 89% task completion rate"
    body5.add_paragraph().text = "JAWS (Windows): 91% task completion rate"
    body5.add_paragraph().text = "VoiceOver (macOS): 85% task completion rate"
    body5.add_paragraph().text = "TalkBack (Android): 78% task completion rate"
    body5.add_paragraph().text = ""
    body5.add_paragraph().text = "Key finding: Tab order inconsistencies on 3 product pages cause navigation confusion for keyboard-only users."

    # ---- Slide 6: Remediation Timeline ----
    sl6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(sl6, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                "Remediation Timeline", font_size=28, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x6B))
    add_textbox(sl6, Inches(0.5), Inches(1.5), Inches(12), Inches(4.5),
                "Phase 1 (Apr 2026): Fix all Level A critical failures\n"
                "  - Add missing alt text to 47 images\n"
                "  - Implement form labels on checkout flow\n"
                "  - Add ARIA landmarks to all page templates\n\n"
                "Phase 2 (May-Jun 2026): Address Level AA issues\n"
                "  - Ensure 4.5:1 color contrast throughout\n"
                "  - Add visible focus indicators to custom components\n"
                "  - Implement skip navigation links\n\n"
                "Phase 3 (Jul-Aug 2026): Enhanced accessibility\n"
                "  - Add keyboard shortcuts for frequent actions\n"
                "  - Implement live regions for dynamic content\n"
                "  - Complete VPAT documentation",
                font_size=14)

    # ---- Slide 7: Training Updates ----
    sl7 = prs.slides.add_slide(prs.slide_layouts[1])
    sl7.shapes.title.text = "Training & Certification"
    body7 = sl7.placeholders[1].text_frame
    body7.text = "32 developers completed IAAP WAS certification"
    body7.add_paragraph().text = "15 designers enrolled in inclusive design workshop"
    body7.add_paragraph().text = "QA team trained on axe-core automated testing"
    body7.add_paragraph().text = "Monthly accessibility office hours: avg 18 attendees"

    # ---- Slide 8: Focus — Tab Order Accessibility (NON-SEQUENTIAL INSERTION) ----
    # This is the key slide. We add shapes in WRONG order so tab order is incorrect.
    # Correct order should be: Title, Left, Center, Right, Footer
    # We add them as: Right, Footer, Center, Title, Left — making tab order wrong.
    sl8 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    # Shape 1 (added first → tab pos 1): RIGHT content box
    add_textbox(sl8, Inches(9.0), Inches(1.8), Inches(3.8), Inches(4.2),
                "Keyboard Navigation\n\n"
                "All interactive elements must be\n"
                "reachable via Tab key. Custom\n"
                "widgets need role='button' or\n"
                "role='link' with proper ARIA\n"
                "attributes. Focus order should\n"
                "follow visual reading flow.",
                font_size=13)

    # Shape 2 (added second → tab pos 2): FOOTER text box
    add_textbox(sl8, Inches(0.5), Inches(6.5), Inches(12), Inches(0.6),
                "Digital Accessibility Team — Q1 2026 Compliance Review — Page 8",
                font_size=10, color=RGBColor(0x80, 0x80, 0x80),
                alignment=PP_ALIGN.CENTER)

    # Shape 3 (added third → tab pos 3): CENTER content box
    add_textbox(sl8, Inches(4.8), Inches(1.8), Inches(3.8), Inches(4.2),
                "Color & Contrast\n\n"
                "Minimum contrast ratio of 4.5:1\n"
                "for normal text and 3:1 for large\n"
                "text (18pt or 14pt bold). Never\n"
                "use color as the sole indicator\n"
                "of meaning — always provide\n"
                "additional visual cues.",
                font_size=13)

    # Shape 4 (added fourth → tab pos 4): TITLE
    add_textbox(sl8, Inches(0.5), Inches(0.3), Inches(12), Inches(1.2),
                "Accessibility Best Practices for Tab Order",
                font_size=28, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x6B))

    # Shape 5 (added fifth → tab pos 5): LEFT content box
    add_textbox(sl8, Inches(0.5), Inches(1.8), Inches(3.8), Inches(4.2),
                "Semantic Structure\n\n"
                "Use proper heading hierarchy\n"
                "(H1 through H6) to organize\n"
                "content. Screen readers use\n"
                "headings for navigation.\n"
                "Landmarks (nav, main, aside)\n"
                "help users jump to sections.",
                font_size=13)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file open
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
