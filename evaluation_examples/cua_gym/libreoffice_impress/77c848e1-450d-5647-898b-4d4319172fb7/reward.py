"""
Reward Script: Verify tab order of shapes on slide 8 of complex_deck.pptx
Task ID: impress_gf5_025
Domain: libreoffice_impress

Tab order in .pptx maps to the order of shape elements in the slide's spTree XML.
python-pptx iterates shapes in spTree order.

Expected golden tab order (positions 0-5):
  0: Title placeholder (empty)
  1: Title text box ("Accessibility Best Practices for Tab Order")
  2: Left content ("Semantic Structure")
  3: Center content ("Color & Contrast")
  4: Right content ("Keyboard Navigation")
  5: Footer ("Digital Accessibility Team...")

Initial (wrong) order:
  0: title_placeholder, 1: right_content, 2: footer, 3: center_content, 4: title_text, 5: left_content

Scoring:
  Component 1 (0.35): Title text is first AND left content is second (positions 1-2)
  Component 2 (0.35): Center content is third AND right content is fourth (positions 3-4)
  Component 3 (0.30): Footer is last AND title placeholder is first (full bookend check)

All components FAIL on initial because at least one sub-condition in each fails.
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_025'


def persist_app_state(domain: str):
    """Save any unsaved GUI state before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(1.0)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def classify_shape(shape):
    """
    Classify a shape into one of the expected categories based on its content and position.
    Returns: 'title_placeholder', 'title_text', 'left_content', 'center_content',
             'right_content', 'footer', or 'unknown'
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    text = ""
    if hasattr(shape, 'text'):
        text = shape.text.strip()

    # Empty placeholder = title placeholder
    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and not text:
        return 'title_placeholder'

    # Footer: positioned at bottom of slide (top > 5000000 EMU)
    if shape.top > 5000000:
        return 'footer'

    # Title text: full-width shape near top with large width
    if shape.top < 500000 and shape.width > 9000000:
        return 'title_text'

    # Content boxes: three boxes at similar vertical position (~1645920 EMU)
    # Distinguish by horizontal position (left coordinate)
    if 1000000 < shape.top < 3000000:
        if shape.left < 2000000:
            return 'left_content'
        elif shape.left < 6000000:
            return 'center_content'
        else:
            return 'right_content'

    return 'unknown'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 8 slides
    if len(prs.slides) < 8:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 8")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[7]  # Slide 8, 0-indexed

    # Precondition: must have exactly 6 shapes on slide 8
    if len(slide.shapes) != 6:
        print(f"FAIL: Slide 8 has {len(slide.shapes)} shapes, expected 6")
        print("REWARD: 0.0")
        return 0.0

    # Classify each shape by its content/position
    shape_classes = []
    for i, shape in enumerate(slide.shapes):
        cls = classify_shape(shape)
        text_preview = shape.text[:60] if hasattr(shape, 'text') else 'N/A'
        shape_classes.append(cls)
        print(f"  Shape {i}: class={cls}, left={shape.left}, top={shape.top}, text={repr(text_preview)}")

    print(f"\nClassified order: {shape_classes}")

    # Component 1: Title text at position 1 AND Left content at position 2 (0.35 points)
    # Initial: pos 1=right_content, pos 2=footer -> FAILS
    # Golden:  pos 1=title_text, pos 2=left_content -> PASSES
    try:
        pos1_ok = (shape_classes[1] == 'title_text')
        pos2_ok = (shape_classes[2] == 'left_content')
        if pos1_ok and pos2_ok:
            print(f"PASS: Component 1 — Title text at pos 1, left content at pos 2 (0.35 pts)")
            total_score += 0.35
        else:
            details = []
            if not pos1_ok:
                details.append(f"pos 1: expected title_text, found {shape_classes[1]}")
            if not pos2_ok:
                details.append(f"pos 2: expected left_content, found {shape_classes[2]}")
            print(f"FAIL: Component 1 — {'; '.join(details)}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Center content at position 3 AND Right content at position 4 (0.35 points)
    # Initial: pos 3=center_content (correct), pos 4=title_text (wrong) -> FAILS (AND gate)
    # Golden:  pos 3=center_content, pos 4=right_content -> PASSES
    try:
        pos3_ok = (shape_classes[3] == 'center_content')
        pos4_ok = (shape_classes[4] == 'right_content')
        if pos3_ok and pos4_ok:
            print(f"PASS: Component 2 — Center content at pos 3, right content at pos 4 (0.35 pts)")
            total_score += 0.35
        else:
            details = []
            if not pos3_ok:
                details.append(f"pos 3: expected center_content, found {shape_classes[3]}")
            if not pos4_ok:
                details.append(f"pos 4: expected right_content, found {shape_classes[4]}")
            print(f"FAIL: Component 2 — {'; '.join(details)}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Footer at position 5 (last) AND title_text not at position 4 or 5 (0.30 points)
    # This ensures the complete reorder: footer moved to end, title moved to front.
    # Initial: pos 5=left_content (not footer) -> FAILS
    # Golden:  pos 5=footer, pos 1=title_text -> PASSES
    try:
        footer_last = (shape_classes[5] == 'footer')
        title_early = (shape_classes[1] == 'title_text')
        if footer_last and title_early:
            print(f"PASS: Component 3 — Footer at last position, title at front (0.30 pts)")
            total_score += 0.30
        else:
            details = []
            if not footer_last:
                details.append(f"pos 5: expected footer, found {shape_classes[5]}")
            if not title_early:
                details.append(f"pos 1: expected title_text, found {shape_classes[1]}")
            print(f"FAIL: Component 3 — {'; '.join(details)}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
