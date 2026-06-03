"""
Initial Setup: Presentation with uniform 18pt body text on master slide
Task ID: impress_ma_035
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_035'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_master_body_uniform_18pt(prs):
    """Set the slide master body placeholder to uniform 18pt regular for all levels."""
    master = prs.slide_masters[0]
    # Find the body placeholder in the slide master
    body_ph = None
    for shape in master.shapes:
        if shape.is_placeholder:
            ph_idx = shape.placeholder_format.idx
            # idx 1 is typically the body/content placeholder on the master
            if ph_idx == 1:
                body_ph = shape
                break

    if body_ph is None:
        # Try to find by looking at the XML for body type
        for shape in master.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                # type 2 = BODY
                if ph_type is not None and ph_type == 2:
                    body_ph = shape
                    break

    if body_ph is not None and body_ph.has_text_frame:
        tf = body_ph.text_frame
        # Set uniform 18pt regular for all existing paragraphs (levels)
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(18)
                run.font.bold = False
                run.font.italic = False
                run.font.name = "Calibri"

    # Also set via XML on the lstStyle of the master body placeholder
    # to ensure all outline levels inherit uniform 18pt
    _set_master_body_lstStyle_uniform(master, 18)


def _set_master_body_lstStyle_uniform(master, size_pt):
    """Set the slide master's body text lstStyle to uniform size for levels 1-9."""
    # Access the slide master XML directly
    sldMaster = master._element
    ns = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    }

    # Find the txStyles element which contains bodyStyle
    txStyles = sldMaster.find('.//p:txStyles', ns)
    if txStyles is not None:
        bodyStyle = txStyles.find('p:bodyStyle', ns)
        if bodyStyle is not None:
            size_hundredths = size_pt * 100  # font size in hundredths of a point
            # Set all levels (lvl1pPr through lvl9pPr) to uniform size
            for lvl in range(1, 10):
                tag = f'a:lvl{lvl}pPr'
                lvlPr = bodyStyle.find(tag, ns)
                if lvlPr is not None:
                    # Find or create defRPr
                    defRPr = lvlPr.find('a:defRPr', ns)
                    if defRPr is not None:
                        defRPr.set('sz', str(int(size_hundredths)))
                        defRPr.set('b', '0')
                        defRPr.set('i', '0')
                    else:
                        defRPr = lvlPr.makeelement(qn('a:defRPr'), {
                            'sz': str(int(size_hundredths)),
                            'b': '0',
                            'i': '0',
                        })
                        lvlPr.append(defRPr)


def create_initial():
    prs = Presentation()

    # Set uniform 18pt on master body text
    set_master_body_uniform_18pt(prs)

    # Lecture slide content - "Outline Lecture" theme
    slide_content = [
        {
            "layout": 0,  # Title slide
            "title": "Advanced Data Analytics",
            "subtitle": "A Comprehensive Lecture Series\nDr. Elena Rodriguez\nSpring 2025"
        },
        {
            "layout": 1,  # Title + Content
            "title": "Course Overview",
            "body": [
                (0, "Introduction to Data Analytics Fundamentals"),
                (1, "Statistical methods and probability theory"),
                (1, "Data collection and preprocessing techniques"),
                (2, "Handling missing values and outliers"),
                (2, "Feature engineering best practices"),
                (0, "Machine Learning Applications"),
                (1, "Supervised and unsupervised learning paradigms"),
            ]
        },
        {
            "layout": 1,
            "title": "Module 1: Statistical Foundations",
            "body": [
                (0, "Descriptive Statistics"),
                (1, "Measures of central tendency"),
                (2, "Mean, median, and mode calculations"),
                (1, "Measures of dispersion"),
                (2, "Variance, standard deviation, and IQR"),
                (0, "Inferential Statistics"),
                (1, "Hypothesis testing frameworks"),
                (2, "t-tests, chi-square, and ANOVA"),
            ]
        },
        {
            "layout": 1,
            "title": "Module 2: Data Preprocessing",
            "body": [
                (0, "Data Cleaning Pipeline"),
                (1, "Identifying data quality issues"),
                (1, "Normalization and standardization"),
                (2, "Min-max scaling techniques"),
                (2, "Z-score normalization approaches"),
                (0, "Feature Selection"),
                (1, "Correlation analysis methods"),
            ]
        },
        {
            "layout": 1,
            "title": "Module 3: Regression Analysis",
            "body": [
                (0, "Linear Regression"),
                (1, "Simple and multiple regression models"),
                (2, "Ordinary least squares estimation"),
                (0, "Logistic Regression"),
                (1, "Binary classification applications"),
                (2, "Maximum likelihood estimation"),
                (1, "Model evaluation metrics"),
            ]
        },
        {
            "layout": 1,
            "title": "Module 4: Classification Techniques",
            "body": [
                (0, "Decision Trees and Random Forests"),
                (1, "Information gain and Gini impurity"),
                (2, "Pruning strategies for overfitting"),
                (0, "Support Vector Machines"),
                (1, "Kernel functions and hyperplanes"),
                (2, "Soft margin classification"),
            ]
        },
        {
            "layout": 1,
            "title": "Module 5: Clustering Methods",
            "body": [
                (0, "Partitioning Approaches"),
                (1, "K-means clustering algorithm"),
                (2, "Elbow method for optimal K selection"),
                (0, "Hierarchical Clustering"),
                (1, "Agglomerative and divisive strategies"),
                (2, "Dendrogram interpretation and analysis"),
                (1, "Distance metrics comparison"),
            ]
        },
        {
            "layout": 1,
            "title": "Module 6: Dimensionality Reduction",
            "body": [
                (0, "Principal Component Analysis"),
                (1, "Eigenvalue decomposition"),
                (2, "Variance explained ratio"),
                (0, "t-SNE and UMAP Visualization"),
                (1, "Non-linear dimensionality reduction"),
                (2, "Perplexity parameter tuning"),
            ]
        },
        {
            "layout": 1,
            "title": "Module 7: Neural Networks",
            "body": [
                (0, "Feedforward Networks"),
                (1, "Activation functions and layers"),
                (2, "ReLU, sigmoid, and tanh characteristics"),
                (0, "Training Optimization"),
                (1, "Gradient descent variants"),
                (2, "Adam, SGD, and RMSprop comparisons"),
                (1, "Regularization techniques"),
            ]
        },
        {
            "layout": 1,
            "title": "Module 8: Deep Learning Applications",
            "body": [
                (0, "Convolutional Neural Networks"),
                (1, "Image recognition and feature extraction"),
                (2, "Pooling layers and stride configuration"),
                (0, "Recurrent Neural Networks"),
                (1, "Sequence modeling and time series"),
                (2, "LSTM and GRU architectures"),
            ]
        },
        {
            "layout": 1,
            "title": "Module 9: Natural Language Processing",
            "body": [
                (0, "Text Preprocessing"),
                (1, "Tokenization and stemming"),
                (2, "Stop word removal strategies"),
                (0, "Language Models"),
                (1, "Word embeddings and transformers"),
                (2, "Attention mechanism fundamentals"),
                (1, "Sentiment analysis applications"),
            ]
        },
        {
            "layout": 1,
            "title": "Module 10: Model Evaluation",
            "body": [
                (0, "Performance Metrics"),
                (1, "Accuracy, precision, recall, and F1-score"),
                (2, "Confusion matrix interpretation"),
                (0, "Cross-Validation"),
                (1, "K-fold and stratified validation"),
                (2, "Leave-one-out cross-validation"),
            ]
        },
        {
            "layout": 1,
            "title": "Research Project Guidelines",
            "body": [
                (0, "Project Requirements"),
                (1, "Select a real-world dataset with 1000+ records"),
                (2, "Kaggle or UCI repository recommended"),
                (1, "Apply minimum three analytical techniques"),
                (0, "Deliverables"),
                (1, "Written report with methodology section"),
                (2, "Visualizations and statistical analysis"),
                (1, "Presentation of findings to class"),
            ]
        },
        {
            "layout": 1,
            "title": "Assessment Schedule",
            "body": [
                (0, "Midterm Examination - March 15, 2025"),
                (1, "Covers Modules 1 through 5"),
                (2, "Multiple choice and short answer format"),
                (0, "Final Examination - May 20, 2025"),
                (1, "Comprehensive coverage of all modules"),
                (2, "Includes practical coding component"),
                (1, "Open-book with approved reference materials"),
            ]
        },
        {
            "layout": 1,
            "title": "Recommended Resources",
            "body": [
                (0, "Textbooks"),
                (1, "Introduction to Statistical Learning by James et al."),
                (2, "Free online edition available"),
                (0, "Online Platforms"),
                (1, "Coursera Machine Learning specialization"),
                (1, "DataCamp Python for Data Science track"),
                (2, "Student accounts provided by department"),
            ]
        },
    ]

    for i, slide_data in enumerate(slide_content):
        layout_idx = slide_data["layout"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data["title"]

        if layout_idx == 0:
            # Title slide - set subtitle
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_data.get("subtitle", "")
        elif "body" in slide_data:
            # Content slide - find body placeholder
            body_ph = None
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    body_ph = ph
                    break

            if body_ph is not None:
                tf = body_ph.text_frame
                tf.clear()
                for j, (level, text) in enumerate(slide_data["body"]):
                    if j == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = text
                    p.level = level

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
