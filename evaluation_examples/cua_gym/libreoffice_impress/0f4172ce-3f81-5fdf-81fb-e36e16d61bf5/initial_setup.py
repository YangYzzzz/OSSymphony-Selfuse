"""
Initial Setup: Duplicate pentagon shape on slide 5 with repositioning and recoloring
Task ID: impress_ndo_063
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_063'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Cm(33.867)  # standard 16:9
    prs.slide_height = Cm(19.05)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Geometric Shapes Portfolio"
    slide1.placeholders[1].text = "Visual Design Department - Q2 2025 Review"

    # --- Slide 2: Content slide with overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide2.shapes.add_textbox(Cm(2), Cm(1), Cm(20), Cm(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Shape Catalog Overview"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    p2 = tf.add_paragraph()
    p2.text = ("This portfolio showcases the primary geometric shapes used "
               "across our branding materials and presentation templates.")
    p2.alignment = PP_ALIGN.LEFT
    for r in p2.runs:
        r.font.size = Pt(14)
        r.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)

    # Add a rectangle shape
    rect = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(8), Cm(8), Cm(6), Cm(4)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x27, 0xAE, 0x60)
    rect.text_frame.paragraphs[0].text = "Rectangle"
    for r in rect.text_frame.paragraphs[0].runs:
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        r.font.size = Pt(14)

    # --- Slide 3: Circles and Ellipses ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox3 = slide3.shapes.add_textbox(Cm(2), Cm(1), Cm(20), Cm(2))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Circles and Ellipses"
    run3 = p3.runs[0]
    run3.font.size = Pt(24)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    circle1 = slide3.shapes.add_shape(
        MSO_SHAPE.OVAL, Cm(4), Cm(5), Cm(5), Cm(5)
    )
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = RGBColor(0xF3, 0x9C, 0x12)

    circle2 = slide3.shapes.add_shape(
        MSO_SHAPE.OVAL, Cm(14), Cm(5), Cm(7), Cm(5)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(0x9B, 0x59, 0xB6)

    # --- Slide 4: Triangles ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox4 = slide4.shapes.add_textbox(Cm(2), Cm(1), Cm(20), Cm(2))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Triangle Variations"
    run4 = p4.runs[0]
    run4.font.size = Pt(24)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    tri1 = slide4.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Cm(5), Cm(6), Cm(6), Cm(6)
    )
    tri1.fill.solid()
    tri1.fill.fore_color.rgb = RGBColor(0xE6, 0x7E, 0x22)

    tri2 = slide4.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Cm(16), Cm(6), Cm(6), Cm(6)
    )
    tri2.fill.solid()
    tri2.fill.fore_color.rgb = RGBColor(0x1A, 0xBC, 0x9C)

    # --- Slide 5: Pentagon (the task slide) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox5 = slide5.shapes.add_textbox(Cm(2), Cm(1), Cm(20), Cm(2))
    tf5 = txBox5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Pentagon Showcase"
    run5 = p5.runs[0]
    run5.font.size = Pt(24)
    run5.font.bold = True
    run5.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # One pentagon at (8cm, 6cm), size 5cm x 5cm, fill #3498DB
    pentagon = slide5.shapes.add_shape(
        MSO_SHAPE.PENTAGON, Cm(8), Cm(6), Cm(5), Cm(5)
    )
    pentagon.fill.solid()
    pentagon.fill.fore_color.rgb = RGBColor(0x34, 0x98, 0xDB)
    pentagon.line.fill.background()  # no outline

    # --- Slide 6: Summary ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox6 = slide6.shapes.add_textbox(Cm(2), Cm(2), Cm(28), Cm(4))
    tf6 = txBox6.text_frame
    tf6.word_wrap = True
    p6 = tf6.paragraphs[0]
    p6.text = "Shape Usage Summary"
    run6 = p6.runs[0]
    run6.font.size = Pt(28)
    run6.font.bold = True
    run6.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    p6b = tf6.add_paragraph()
    p6b.text = ("Our design system leverages geometric primitives to create "
                "clean, recognizable visual patterns. Each shape serves a "
                "distinct purpose in our presentation hierarchy.")
    for r in p6b.runs:
        r.font.size = Pt(14)
        r.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
