"""
Reward Script: Duplicate pentagon on slide 5, position copy 2cm right, change fill to #E74C3C
Task ID: impress_ndo_063
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): Two pentagon shapes on slide 5
  Component 2 (0.3): Original pentagon preserved at (8cm, 6cm) with fill #3498DB
  Component 3 (0.2): Copy pentagon at ~(15cm, 6cm) with size 5x5cm
  Component 4 (0.2): Copy pentagon fill color is #E74C3C
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_063'


def persist_app_state(domain):
    """Save any unsaved changes in LibreOffice."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def is_approx(val1_cm, val2_cm, tolerance_cm=0.5):
    """Check if two values in cm are approximately equal."""
    return abs(val1_cm - val2_cm) <= tolerance_cm


def emu_to_cm(emu):
    """Convert EMU to centimeters."""
    return emu / 360000.0


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # slide 5 (0-indexed)

    # Find all pentagon shapes on slide 5
    # Pentagon shapes have prstGeom 'homePlate' or auto_shape_type PENTAGON
    pentagons = []
    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                el = shape._element
                ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
                prstGeom = el.find(f'.//{ns}prstGeom')
                if prstGeom is not None and prstGeom.get('prst') == 'homePlate':
                    pentagons.append(shape)
        except Exception as e:
            print(f"WARN: Error inspecting shape: {e}")

    print(f"INFO: Found {len(pentagons)} pentagon shapes on slide 5")

    # Component 1: Two pentagon shapes exist on slide 5 (0.3 points)
    # This FAILS on initial (1 pentagon) and PASSES on golden (2 pentagons)
    try:
        if len(pentagons) >= 2:
            print(f"PASS: Component 1 - Found {len(pentagons)} pentagons on slide 5 (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 - Expected 2 pentagons, found {len(pentagons)}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # If fewer than 2 pentagons, cannot verify further components meaningfully
    if len(pentagons) < 2:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Categorize pentagons: original (fill #3498DB) and copy (fill #E74C3C)
    original = None
    copy_shape = None
    for p in pentagons:
        try:
            fill = p.fill
            if fill.type is not None:
                color_hex = str(fill.fore_color.rgb).upper()
                if color_hex == '3498DB':
                    original = p
                elif color_hex == 'E74C3C':
                    copy_shape = p
        except Exception as e:
            print(f"WARN: Could not read fill color for pentagon: {e}")

    # Component 2: Original pentagon preserved at (8cm, 6cm), size 5x5cm, fill #3498DB (0.3 points)
    # This FAILS on initial because there's only 1 pentagon (component 1 gates this)
    # Actually, we need this to only score when there are 2 pentagons AND original is intact
    try:
        if original is not None:
            left_cm = emu_to_cm(original.left)
            top_cm = emu_to_cm(original.top)
            w_cm = emu_to_cm(original.width)
            h_cm = emu_to_cm(original.height)

            pos_ok = is_approx(left_cm, 8.0) and is_approx(top_cm, 6.0)
            size_ok = is_approx(w_cm, 5.0) and is_approx(h_cm, 5.0)

            if pos_ok and size_ok:
                print(f"PASS: Component 2 - Original pentagon at ({left_cm:.2f}cm, {top_cm:.2f}cm), "
                      f"size ({w_cm:.2f}cm x {h_cm:.2f}cm), fill #3498DB (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 2 - Original pentagon position/size mismatch: "
                      f"pos=({left_cm:.2f}, {top_cm:.2f}), size=({w_cm:.2f}, {h_cm:.2f})")
        else:
            print(f"FAIL: Component 2 - No pentagon with fill #3498DB found")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Copy pentagon at ~(15cm, 6cm) with size 5x5cm (0.2 points)
    try:
        if copy_shape is not None:
            left_cm = emu_to_cm(copy_shape.left)
            top_cm = emu_to_cm(copy_shape.top)
            w_cm = emu_to_cm(copy_shape.width)
            h_cm = emu_to_cm(copy_shape.height)

            # Task says "2cm to the right of the original". Original ends at 8+5=13cm, so copy at 15cm
            pos_ok = is_approx(left_cm, 15.0, tolerance_cm=1.0) and is_approx(top_cm, 6.0)
            size_ok = is_approx(w_cm, 5.0) and is_approx(h_cm, 5.0)

            if pos_ok and size_ok:
                print(f"PASS: Component 3 - Copy pentagon at ({left_cm:.2f}cm, {top_cm:.2f}cm), "
                      f"size ({w_cm:.2f}cm x {h_cm:.2f}cm) (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 3 - Copy pentagon position/size mismatch: "
                      f"pos=({left_cm:.2f}, {top_cm:.2f}), size=({w_cm:.2f}, {h_cm:.2f})")
        else:
            print(f"FAIL: Component 3 - No pentagon with fill #E74C3C found")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Copy pentagon fill color is #E74C3C (0.2 points)
    try:
        if copy_shape is not None:
            fill = copy_shape.fill
            color_hex = str(fill.fore_color.rgb).upper()
            if color_hex == 'E74C3C':
                print(f"PASS: Component 4 - Copy pentagon fill is #{color_hex} (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 4 - Copy fill expected #E74C3C, found #{color_hex}")
        else:
            print(f"FAIL: Component 4 - No copy pentagon found to check fill color")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
