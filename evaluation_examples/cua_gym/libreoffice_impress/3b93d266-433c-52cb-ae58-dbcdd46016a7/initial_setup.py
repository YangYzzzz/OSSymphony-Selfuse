"""
Initial Setup: Academic conference presentation with 5 slides, NO speaker notes.
Task ID: osworld_impress_slide_notes_008
Domain: libreoffice_impress

Creates a 5-slide academic conference presentation WITHOUT any speaker notes.
The agent's task will be to add specific notes to each slide.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_slide_notes_008'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard widescreen slide dimensions
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
    slide1.shapes.title.text = "Neural Architecture Search for Low-Resource NLP Tasks"
    subtitle = slide1.placeholders[1]
    subtitle.text = "Dr. Emily Hartmann, Dr. Raj Patel\nDepartment of Computer Science, Stanford University\nInternational Conference on Machine Learning (ICML 2025)"

    # --- Slide 2: Introduction / Background ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content layout
    slide2.shapes.title.text = "Research Problem & Motivation"
    body2 = slide2.placeholders[1]
    tf2 = body2.text_frame
    tf2.text = "Challenges in Low-Resource NLP"
    p2_1 = tf2.add_paragraph()
    p2_1.text = "Limited labeled data for minority languages and specialized domains"
    p2_1.level = 1
    p2_2 = tf2.add_paragraph()
    p2_2.text = "Manual architecture design is expensive and suboptimal"
    p2_2.level = 1
    p2_3 = tf2.add_paragraph()
    p2_3.text = "Existing AutoML approaches ignore resource constraints"
    p2_3.level = 1
    p2_4 = tf2.add_paragraph()
    p2_4.text = "Gap: No principled method for resource-aware NAS in NLP"
    p2_4.level = 1
    p2_5 = tf2.add_paragraph()
    p2_5.text = "Over 7,000 languages worldwide lack adequate NLP tooling"
    p2_5.level = 2

    # --- Slide 3: Methodology ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content layout
    slide3.shapes.title.text = "Methodology & Data Sources"
    body3 = slide3.placeholders[1]
    tf3 = body3.text_frame
    tf3.text = "Proposed Approach: Resource-Aware NAS (RA-NAS)"
    p3_1 = tf3.add_paragraph()
    p3_1.text = "Differentiable architecture search with resource penalty"
    p3_1.level = 1
    p3_2 = tf3.add_paragraph()
    p3_2.text = "Multi-objective optimization: accuracy vs. FLOPs tradeoff"
    p3_2.level = 1
    p3_3 = tf3.add_paragraph()
    p3_3.text = "Data Sources"
    p3_3.level = 0
    p3_4 = tf3.add_paragraph()
    p3_4.text = "Universal Dependencies Treebank (93 languages, 250K sentences)"
    p3_4.level = 1
    p3_5 = tf3.add_paragraph()
    p3_5.text = "XTREME multilingual benchmark (9 tasks, 40 languages)"
    p3_5.level = 1
    p3_6 = tf3.add_paragraph()
    p3_6.text = "Internal medical domain corpus (12K annotated records)"
    p3_6.level = 1

    # --- Slide 4: Results ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content layout
    slide4.shapes.title.text = "Key Experimental Results"
    body4 = slide4.placeholders[1]
    tf4 = body4.text_frame
    tf4.text = "Performance Comparison"
    p4_1 = tf4.add_paragraph()
    p4_1.text = "RA-NAS achieves 83.7% F1 on NER (vs. 79.2% BERT baseline)"
    p4_1.level = 1
    p4_2 = tf4.add_paragraph()
    p4_2.text = "3.4× reduction in inference FLOPs with <1% accuracy drop"
    p4_2.level = 1
    p4_3 = tf4.add_paragraph()
    p4_3.text = "Generalizes across 12 language families"
    p4_3.level = 1
    p4_4 = tf4.add_paragraph()
    p4_4.text = "Ablation Studies"
    p4_4.level = 0
    p4_5 = tf4.add_paragraph()
    p4_5.text = "Resource penalty coefficient α = 0.15 yields best Pareto front"
    p4_5.level = 1
    p4_6 = tf4.add_paragraph()
    p4_6.text = "Transfer from high-resource languages improves low-resource by 6.2%"
    p4_6.level = 1

    # --- Slide 5: Discussion ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content layout
    slide5.shapes.title.text = "Implications, Limitations & Future Work"
    body5 = slide5.placeholders[1]
    tf5 = body5.text_frame
    tf5.text = "Contributions"
    p5_1 = tf5.add_paragraph()
    p5_1.text = "First resource-aware NAS framework for low-resource NLP"
    p5_1.level = 1
    p5_2 = tf5.add_paragraph()
    p5_2.text = "Open-source toolkit: github.com/hartmann-lab/ra-nas"
    p5_2.level = 1
    p5_3 = tf5.add_paragraph()
    p5_3.text = "Limitations"
    p5_3.level = 0
    p5_4 = tf5.add_paragraph()
    p5_4.text = "Search cost still requires 8× V100 GPUs for 48 hours"
    p5_4.level = 1
    p5_5 = tf5.add_paragraph()
    p5_5.text = "Future Work"
    p5_5.level = 0
    p5_6 = tf5.add_paragraph()
    p5_6.text = "Zero-shot transfer to unseen language families"
    p5_6.level = 1
    p5_7 = tf5.add_paragraph()
    p5_7.text = "Integration with continual learning for dynamic adaptation"
    p5_7.level = 1

    # IMPORTANT: No speaker notes are added to any slide (task requirement)
    # Verify notes are empty
    for i, slide in enumerate(prs.slides):
        notes = slide.notes_slide.notes_text_frame.text.strip()
        assert notes == "", f"Slide {i+1} should have no notes, found: {notes!r}"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
