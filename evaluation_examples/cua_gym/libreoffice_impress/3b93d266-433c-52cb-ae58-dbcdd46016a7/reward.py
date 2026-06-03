"""
Reward Script: Add speaker notes to all 5 slides of a conference talk presentation
Task ID: osworld_impress_slide_notes_008
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 1 note = 'Introduce research problem and motivation' (0.2 pts)
  Component 2: Slide 2 note = 'Explain methodology and data sources' (0.2 pts)
  Component 3: Slide 3 note = 'Present key experimental results' (0.2 pts)
  Component 4: Slide 4 note = 'Discuss implications and limitations' (0.2 pts)
  Component 5: Slide 5 note = 'Summarize contributions and future work' (0.2 pts)
  Total: 1.0
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_slide_notes_008'

# Expected notes per slide (1-indexed descriptions, stored as 0-indexed list)
EXPECTED_NOTES = [
    'Introduce research problem and motivation',   # slide 1
    'Explain methodology and data sources',        # slide 2
    'Present key experimental results',            # slide 3
    'Discuss implications and limitations',        # slide 4
    'Summarize contributions and future work',     # slide 5
]


def get_slide_notes(slide):
    """Return stripped note text from a slide, empty string if none."""
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def verify_task(file_path):
    """
    Verify that each of the 5 slides has the specified speaker note added.
    Each slide is worth 0.2 points; partial credit is awarded per slide.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation — gate check
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Verify the presentation has 5 slides (precondition gate)
    num_slides = len(prs.slides)
    if num_slides != 5:
        print(f"CRITICAL: Expected 5 slides, found {num_slides}. Aborting.")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide 1 note = 'Introduce research problem and motivation' (0.2 pts)
    try:
        actual_note = get_slide_notes(prs.slides[0])
        expected_note = EXPECTED_NOTES[0]
        if actual_note == expected_note:
            print(f"PASS: Component 1 — Slide 1 note matches expected (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 1 — Slide 1 note: expected {repr(expected_note)}, found {repr(actual_note)}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 2 note = 'Explain methodology and data sources' (0.2 pts)
    try:
        actual_note = get_slide_notes(prs.slides[1])
        expected_note = EXPECTED_NOTES[1]
        if actual_note == expected_note:
            print(f"PASS: Component 2 — Slide 2 note matches expected (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 2 — Slide 2 note: expected {repr(expected_note)}, found {repr(actual_note)}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 3 note = 'Present key experimental results' (0.2 pts)
    try:
        actual_note = get_slide_notes(prs.slides[2])
        expected_note = EXPECTED_NOTES[2]
        if actual_note == expected_note:
            print(f"PASS: Component 3 — Slide 3 note matches expected (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 3 — Slide 3 note: expected {repr(expected_note)}, found {repr(actual_note)}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 4 note = 'Discuss implications and limitations' (0.2 pts)
    try:
        actual_note = get_slide_notes(prs.slides[3])
        expected_note = EXPECTED_NOTES[3]
        if actual_note == expected_note:
            print(f"PASS: Component 4 — Slide 4 note matches expected (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 4 — Slide 4 note: expected {repr(expected_note)}, found {repr(actual_note)}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slide 5 note = 'Summarize contributions and future work' (0.2 pts)
    try:
        actual_note = get_slide_notes(prs.slides[4])
        expected_note = EXPECTED_NOTES[4]
        if actual_note == expected_note:
            print(f"PASS: Component 5 — Slide 5 note matches expected (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 5 — Slide 5 note: expected {repr(expected_note)}, found {repr(actual_note)}")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point: test against the canonical artifact path on the VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
