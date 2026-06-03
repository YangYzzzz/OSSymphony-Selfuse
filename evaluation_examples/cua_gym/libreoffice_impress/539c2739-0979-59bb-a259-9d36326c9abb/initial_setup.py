"""
Initial Setup: Create a 30-slide Training Manual presentation with default layouts
Task ID: impress_gf2_039
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_039'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Training manual content for 30 slides
    chapters = [
        ("Chapter 1: Introduction to Safety Protocols", [
            ("Welcome & Course Overview",
             "This training manual covers essential workplace safety protocols.\n"
             "Duration: 8 hours | Certification: OSHA Level 2\n"
             "Instructor: Dr. Rebecca Torres, Safety Engineering Lead"),
            ("Learning Objectives",
             "By the end of this module, participants will be able to:\n"
             "- Identify common workplace hazards\n"
             "- Apply correct PPE selection procedures\n"
             "- Execute emergency evacuation protocols\n"
             "- Document safety incidents properly"),
            ("Course Agenda - Day 1",
             "09:00 - 09:30  Registration & Introduction\n"
             "09:30 - 10:45  Module 1: Hazard Identification\n"
             "10:45 - 11:00  Break\n"
             "11:00 - 12:30  Module 2: PPE Requirements\n"
             "12:30 - 13:30  Lunch\n"
             "13:30 - 15:00  Module 3: Emergency Procedures"),
            ("Pre-Assessment Questions",
             "1. What is the primary purpose of a Job Safety Analysis?\n"
             "2. Name three types of respiratory protection equipment.\n"
             "3. What color designates a fire exit sign in ISO standards?\n"
             "4. Define 'near miss' in the context of safety reporting."),
        ]),
        ("Chapter 2: Hazard Identification & Risk Assessment", [
            ("Types of Workplace Hazards",
             "Physical: Noise, vibration, radiation, temperature extremes\n"
             "Chemical: Solvents, dust, fumes, gases\n"
             "Biological: Bacteria, viruses, fungi, allergens\n"
             "Ergonomic: Repetitive motion, awkward postures, heavy lifting\n"
             "Psychosocial: Stress, violence, bullying, fatigue"),
            ("Risk Assessment Matrix",
             "Severity Levels:\n"
             "  1 - Negligible: First aid only\n"
             "  2 - Minor: Medical treatment, no lost time\n"
             "  3 - Moderate: Lost time injury, temporary disability\n"
             "  4 - Major: Permanent disability\n"
             "  5 - Critical: Fatality or multiple casualties"),
            ("Case Study: Chemical Spill at Plant B",
             "On March 15, 2024, a 200L drum of sodium hydroxide (NaOH)\n"
             "was punctured during forklift operations at Warehouse B.\n"
             "Response time: 4 minutes | Affected area: 15m radius\n"
             "Personnel exposed: 3 | Injuries: 1 minor chemical burn\n"
             "Root cause: Improper stacking and missing containment berm"),
            ("Exercise: Conduct a Walkthrough Inspection",
             "Using the checklist on the following page, identify hazards in\n"
             "the provided facility photographs.\n"
             "Time allowed: 20 minutes\n"
             "Work in pairs and document findings on Form SA-102."),
        ]),
        ("Chapter 3: Personal Protective Equipment", [
            ("PPE Selection Guidelines",
             "Head Protection: Hard hats (Class E, G, or C)\n"
             "Eye Protection: Safety glasses, goggles, face shields\n"
             "Hearing Protection: Earplugs (NRR 25-33), earmuffs (NRR 20-31)\n"
             "Hand Protection: Chemical-resistant gloves, cut-resistant gloves\n"
             "Foot Protection: Steel-toe boots, metatarsal guards"),
            ("PPE Inspection & Maintenance Schedule",
             "Daily: Visual inspection before each use\n"
             "Weekly: Detailed check for cracks, wear, degradation\n"
             "Monthly: Full cleaning and sanitization\n"
             "Quarterly: Professional inspection and testing\n"
             "Annual: Replacement per manufacturer guidelines"),
            ("Proper Donning & Doffing Procedures",
             "Order of Donning (putting on):\n"
             "1. Coveralls or gown  2. Boots or shoe covers\n"
             "3. Respirator or mask  4. Goggles or face shield\n"
             "5. Inner gloves  6. Outer gloves\n\n"
             "Order of Doffing (removing) - REVERSE order:\n"
             "Always wash hands between each removal step."),
            ("Quiz: PPE Knowledge Check",
             "1. When should a hard hat be replaced? ___________\n"
             "2. What NRR rating is needed for 95dB environments? ___\n"
             "3. True/False: Contact lenses can be worn with safety goggles.\n"
             "4. Match the glove material to the chemical hazard:\n"
             "   a) Nitrile  b) Neoprene  c) Butyl rubber  d) PVC"),
        ]),
        ("Chapter 4: Emergency Response Procedures", [
            ("Emergency Action Plan Overview",
             "Every facility must maintain a current Emergency Action Plan (EAP)\n"
             "per OSHA 29 CFR 1910.38. Key components include:\n"
             "- Evacuation routes and assembly points\n"
             "- Communication protocols and alarm systems\n"
             "- Emergency contact list and chain of command\n"
             "- Accounting for all personnel after evacuation"),
            ("Fire Response Protocol",
             "Remember R.A.C.E.:\n"
             "R - Rescue anyone in immediate danger\n"
             "A - Activate the fire alarm system\n"
             "C - Contain the fire (close doors/windows)\n"
             "E - Extinguish or Evacuate\n\n"
             "Fire extinguisher types: A (ordinary), B (flammable liquids),\n"
             "C (electrical), D (combustible metals), K (cooking oils)"),
            ("Medical Emergency Response",
             "Step 1: Assess the scene for safety\n"
             "Step 2: Call 911 or facility emergency number (ext. 5555)\n"
             "Step 3: Provide first aid if trained and safe to do so\n"
             "Step 4: Do not move the victim unless immediate danger exists\n"
             "Step 5: Stay with victim until EMS arrives\n"
             "Step 6: Complete Incident Report Form IR-200 within 24 hours"),
            ("Evacuation Drill Checklist",
             "Pre-Drill: Notify security, assign observers, test alarm system\n"
             "During Drill: Record evacuation start time, monitor routes\n"
             "Post-Drill: Headcount at assembly point, record completion time\n"
             "Acceptable evacuation time: < 5 minutes for buildings < 3 floors\n"
             "Target participation rate: 100% of on-site personnel"),
        ]),
        ("Chapter 5: Documentation & Compliance", [
            ("Incident Reporting Requirements",
             "All incidents must be reported within the following timeframes:\n"
             "- Near misses: Within 24 hours to supervisor\n"
             "- First aid cases: Same day, Form IR-100\n"
             "- Lost time injuries: Immediately, Forms IR-200 + IR-201\n"
             "- Fatalities: Immediately to OSHA within 8 hours\n"
             "- Hospitalizations: Within 24 hours to OSHA"),
            ("Regulatory Compliance Summary",
             "OSHA General Duty Clause (Section 5(a)(1)):\n"
             "Employers must provide a workplace free from recognized hazards.\n\n"
             "Key Standards Applicable to This Facility:\n"
             "29 CFR 1910.132 - PPE General Requirements\n"
             "29 CFR 1910.134 - Respiratory Protection\n"
             "29 CFR 1910.147 - Lockout/Tagout\n"
             "29 CFR 1910.1200 - Hazard Communication"),
            ("Training Records Management",
             "Maintain records for each employee including:\n"
             "- Initial orientation date and content covered\n"
             "- Annual refresher training completion dates\n"
             "- Competency assessment results (pass/fail + scores)\n"
             "- Equipment-specific certifications (forklift, crane, etc.)\n"
             "Retention period: Duration of employment + 3 years"),
            ("Final Assessment & Certification",
             "To receive OSHA Level 2 Safety Certification:\n"
             "- Score 80% or higher on written examination\n"
             "- Complete all practical exercises satisfactorily\n"
             "- Demonstrate proper PPE selection and use\n"
             "- Submit completed workbook with all exercises\n\n"
             "Certification valid for: 2 years from date of completion"),
            ("Course Evaluation & Feedback",
             "Please rate each module on a scale of 1-5:\n"
             "Module 1 - Hazard Identification: ___\n"
             "Module 2 - PPE Requirements: ___\n"
             "Module 3 - Emergency Procedures: ___\n"
             "Module 4 - Documentation: ___\n\n"
             "Additional comments: _________________________________"),
            ("Thank You & Contact Information",
             "Thank you for completing this safety training program.\n\n"
             "For questions or follow-up:\n"
             "Dr. Rebecca Torres - rtorres@safetyfirst.com\n"
             "Safety Department Hotline: 1-800-555-SAFE\n"
             "Online Resources: https://safety.internal.com/training\n\n"
             "Remember: Safety is everyone's responsibility!"),
            ("Appendix A: Emergency Contact Numbers",
             "Plant Manager: ext. 1001 - David Nakamura\n"
             "Safety Director: ext. 1015 - Dr. Rebecca Torres\n"
             "Security Office: ext. 1100 (24/7)\n"
             "Medical Clinic: ext. 1200 (Mon-Fri 7am-5pm)\n"
             "Poison Control: 1-800-222-1222\n"
             "Local Fire Department: 555-0190\n"
             "Local Hospital (St. Mary's): 555-0200"),
            ("Appendix B: Abbreviations & Glossary",
             "ALARA - As Low As Reasonably Achievable\n"
             "EAP - Emergency Action Plan\n"
             "JSA - Job Safety Analysis\n"
             "LOTO - Lockout/Tagout\n"
             "MSDS - Material Safety Data Sheet\n"
             "NRR - Noise Reduction Rating\n"
             "PPE - Personal Protective Equipment\n"
             "SDS - Safety Data Sheet (replaced MSDS)"),
            ("Appendix C: Revision History",
             "Version 4.2 - March 2024: Updated PPE standards\n"
             "Version 4.1 - January 2024: Added chemical spill case study\n"
             "Version 4.0 - September 2023: Major revision for OSHA updates\n"
             "Version 3.5 - June 2023: Added ergonomic hazards section\n"
             "Version 3.0 - January 2023: Initial release\n\n"
             "Next scheduled review: September 2024"),
        ]),
    ]

    slide_num = 0
    for chapter_title, slides_content in chapters:
        # Chapter title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
        slide.shapes.title.text = chapter_title
        try:
            slide.placeholders[1].text = "Workplace Safety Training Manual 2024"
        except (KeyError, IndexError):
            pass
        slide_num += 1

        # Content slides
        for title, body in slides_content:
            if slide_num >= 30:
                break
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
            slide.shapes.title.text = title
            try:
                tf = slide.placeholders[1].text_frame
                tf.text = body
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(14)
            except (KeyError, IndexError):
                pass
            slide_num += 1

        if slide_num >= 30:
            break

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
