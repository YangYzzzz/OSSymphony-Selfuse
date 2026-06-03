"""
Initial Setup: 6-slide product demo presentation for export task
Task ID: osworld_impress_export_image_007
Domain: libreoffice_impress

Creates a 6-slide product demo deck at /home/user/osworld_impress_export_image_007.pptx
and ensures ~/slide_exports/ directory exists (empty — no PNG files yet).
Opens the file in LibreOffice Impress for the GUI agent.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_export_image_007'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
EXPORT_DIR = f'{WORKDIR}/slide_exports'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    # Ensure slide_exports directory exists and is EMPTY (no PNG exports yet)
    os.makedirs(EXPORT_DIR, exist_ok=True)
    # Remove any existing PNG files (idempotent)
    for f in os.listdir(EXPORT_DIR):
        if f.endswith('.png'):
            os.remove(os.path.join(EXPORT_DIR, f))

    prs = Presentation()
    # Use standard 16:9 widescreen (default 10x7.5 inches)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # -----------------------------------------------------------------------
    # Slide 1: Title Slide — "NovaTech Pro 2025 — Product Demo"
    # -----------------------------------------------------------------------
    layout_title = prs.slide_layouts[0]  # Title Slide
    slide1 = prs.slides.add_slide(layout_title)
    slide1.shapes.title.text = "NovaTech Pro 2025"
    slide1.placeholders[1].text = "Product Demo — Q1 2025"

    # Background: deep navy
    bg1 = slide1.background.fill
    bg1.solid()
    bg1.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3E)

    tf_title = slide1.shapes.title.text_frame.paragraphs[0]
    for run in tf_title.runs:
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.size = Pt(44)
        run.font.bold = True

    tf_sub = slide1.placeholders[1].text_frame.paragraphs[0]
    for run in tf_sub.runs:
        run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xFF)
        run.font.size = Pt(24)

    # -----------------------------------------------------------------------
    # Slide 2: Agenda / Overview
    # -----------------------------------------------------------------------
    layout_content = prs.slide_layouts[1]  # Title + Content
    slide2 = prs.slides.add_slide(layout_content)
    slide2.shapes.title.text = "Agenda"

    body2 = slide2.placeholders[1]
    tf2 = body2.text_frame
    tf2.text = "1. Company Overview"
    items2 = [
        "2. Product Features",
        "3. Market Opportunity",
        "4. Customer Success Stories",
        "5. Pricing & Plans",
    ]
    for item in items2:
        p = tf2.add_paragraph()
        p.text = item
        p.level = 0
        p.space_before = Pt(6)

    for para in tf2.paragraphs:
        for run in para.runs:
            run.font.size = Pt(20)

    # -----------------------------------------------------------------------
    # Slide 3: Company Overview
    # -----------------------------------------------------------------------
    slide3 = prs.slides.add_slide(layout_content)
    slide3.shapes.title.text = "Company Overview"

    body3 = slide3.placeholders[1]
    tf3 = body3.text_frame
    tf3.text = "Founded in 2018 — San Francisco, CA"
    overview_items = [
        "350+ enterprise customers across 28 countries",
        "Annual Recurring Revenue: $42M (up 87% YoY)",
        "Team of 180 engineers, designers & data scientists",
        "Backed by Sequoia Capital & Andreessen Horowitz",
        "ISO 27001 certified — SOC 2 Type II compliant",
    ]
    for item in overview_items:
        p = tf3.add_paragraph()
        p.text = item
        p.level = 1

    for para in tf3.paragraphs:
        for run in para.runs:
            run.font.size = Pt(18)

    # -----------------------------------------------------------------------
    # Slide 4: Product Features
    # -----------------------------------------------------------------------
    slide4 = prs.slides.add_slide(layout_content)
    slide4.shapes.title.text = "Key Product Features"

    body4 = slide4.placeholders[1]
    tf4 = body4.text_frame
    features = [
        "Real-time Analytics Dashboard — sub-50ms query latency",
        "AI-Powered Anomaly Detection — 99.3% precision",
        "One-click Integrations with Salesforce, HubSpot & Slack",
        "Enterprise SSO (SAML 2.0 / OAuth 2.0)",
        "Automated Compliance Reporting (GDPR, HIPAA, SOX)",
        "Custom API Gateway — 10M requests/day on Standard tier",
    ]
    tf4.text = features[0]
    for feat in features[1:]:
        p = tf4.add_paragraph()
        p.text = feat
        p.level = 0

    for para in tf4.paragraphs:
        for run in para.runs:
            run.font.size = Pt(17)
            run.font.bold = False

    # -----------------------------------------------------------------------
    # Slide 5: Customer Success Stories
    # -----------------------------------------------------------------------
    slide5 = prs.slides.add_slide(layout_content)
    slide5.shapes.title.text = "Customer Success Stories"

    body5 = slide5.placeholders[1]
    tf5 = body5.text_frame
    tf5.text = "GlobalRetail Inc."
    stories = [
        "  Reduced churn by 34% in 6 months using Predictive Signals",
        "FinServe Bank",
        "  Automated 92% of compliance workflows, saving 1,200 hrs/month",
        "MedLogistics GmbH",
        "  Cut operational costs by $2.8M annually via route optimization",
    ]
    for s in stories:
        p = tf5.add_paragraph()
        p.text = s
        p.level = 1 if s.startswith("  ") else 0

    for para in tf5.paragraphs:
        for run in para.runs:
            run.font.size = Pt(17)

    # -----------------------------------------------------------------------
    # Slide 6: Pricing & Call to Action
    # -----------------------------------------------------------------------
    slide6 = prs.slides.add_slide(layout_content)
    slide6.shapes.title.text = "Pricing & Next Steps"

    body6 = slide6.placeholders[1]
    tf6 = body6.text_frame
    tf6.text = "Starter: $499/month — up to 10 users, 5 integrations"
    pricing = [
        "Growth: $1,499/month — up to 50 users, unlimited integrations",
        "Enterprise: Custom pricing — dedicated infrastructure + SLA",
        "",
        "30-day free trial — no credit card required",
        "Contact: sales@novatech.io | +1 (415) 555-0192",
        "Schedule a live demo at novatech.io/demo",
    ]
    for line in pricing:
        p = tf6.add_paragraph()
        p.text = line
        p.level = 0

    for para in tf6.paragraphs:
        for run in para.runs:
            run.font.size = Pt(17)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'slide_exports directory ready (empty): {EXPORT_DIR}')

    # GUI-ready startup: open the presentation in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
