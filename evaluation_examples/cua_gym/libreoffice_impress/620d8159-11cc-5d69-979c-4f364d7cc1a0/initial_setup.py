"""
Initial Setup: Create a 10-slide business presentation and prepare html_slides directory
Task ID: impress_gf5_049
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_049'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
HTML_DIR = f'{WORKDIR}/html_slides'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_slide(slide, left, top, width, height, text, font_size=18,
                       bold=False, alignment=PP_ALIGN.LEFT, color=None):
    """Helper to add a textbox with formatted text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ======= Slide 1: Title Slide =======
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only
    add_text_to_slide(slide1, 1.5, 1.5, 7, 1.5,
                      "Q3 2025 Strategic Planning Review",
                      font_size=32, bold=True, alignment=PP_ALIGN.CENTER,
                      color=(0x1B, 0x3A, 0x6B))
    add_text_to_slide(slide1, 2, 3.5, 6, 1,
                      "Meridian Technologies Inc.",
                      font_size=20, alignment=PP_ALIGN.CENTER,
                      color=(0x4A, 0x4A, 0x4A))
    add_text_to_slide(slide1, 2, 4.5, 6, 0.8,
                      "Prepared by: Sarah Chen, VP of Strategy\nSeptember 15, 2025",
                      font_size=14, alignment=PP_ALIGN.CENTER,
                      color=(0x80, 0x80, 0x80))

    # ======= Slide 2: Executive Summary =======
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_to_slide(slide2, 0.5, 0.3, 9, 0.8,
                      "Executive Summary",
                      font_size=28, bold=True, color=(0x1B, 0x3A, 0x6B))
    summary_text = (
        "Meridian Technologies achieved 18% year-over-year revenue growth in Q3, "
        "reaching $142.3M in total revenue. Our cloud services division led with 34% "
        "growth, while the enterprise solutions segment maintained steady 8% expansion. "
        "Key wins include the Pinnacle Healthcare contract ($12.4M ARR) and expansion "
        "of the DataBridge platform into three new markets. Operating margins improved "
        "to 22.1%, up from 19.8% in Q2, driven by infrastructure optimization and "
        "reduced customer acquisition costs."
    )
    add_text_to_slide(slide2, 0.5, 1.3, 9, 4, summary_text, font_size=16)

    # ======= Slide 3: Revenue Breakdown =======
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_to_slide(slide3, 0.5, 0.3, 9, 0.8,
                      "Revenue Breakdown by Division",
                      font_size=28, bold=True, color=(0x1B, 0x3A, 0x6B))
    # Add a table for revenue data
    rows, cols = 6, 4
    tbl = slide3.shapes.add_table(rows, cols,
                                   Inches(0.8), Inches(1.5),
                                   Inches(8.4), Inches(3.5))
    table = tbl.table
    headers = ["Division", "Q3 Revenue", "Q2 Revenue", "Growth %"]
    data = [
        ["Cloud Services", "$52.8M", "$39.4M", "+34.0%"],
        ["Enterprise Solutions", "$41.6M", "$38.5M", "+8.1%"],
        ["Data Analytics", "$28.3M", "$24.1M", "+17.4%"],
        ["Consulting & Support", "$14.2M", "$13.8M", "+2.9%"],
        ["Other / Licensing", "$5.4M", "$4.7M", "+14.9%"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # ======= Slide 4: Market Position =======
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_to_slide(slide4, 0.5, 0.3, 9, 0.8,
                      "Market Position & Competitive Landscape",
                      font_size=28, bold=True, color=(0x1B, 0x3A, 0x6B))
    market_text = (
        "Market Share: 14.2% (up from 12.8% in Q2)\n\n"
        "Key Competitors:\n"
        "  - TechVista Corp: 22.1% market share (declining)\n"
        "  - NovaSoft Industries: 17.5% market share (stable)\n"
        "  - Quantum Digital: 11.3% market share (growing)\n\n"
        "Our differentiation continues to be the integrated DataBridge platform, "
        "which combines real-time analytics with AI-driven automation. Customer "
        "retention rate stands at 94.7%, the highest in our segment."
    )
    add_text_to_slide(slide4, 0.5, 1.3, 9, 5, market_text, font_size=15)

    # ======= Slide 5: Product Roadmap =======
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_to_slide(slide5, 0.5, 0.3, 9, 0.8,
                      "Product Roadmap — Q4 2025 & Beyond",
                      font_size=28, bold=True, color=(0x1B, 0x3A, 0x6B))
    roadmap_text = (
        "Q4 2025 Deliverables:\n"
        "  1. DataBridge v3.2 — Enhanced ML pipeline integration\n"
        "  2. CloudGuard Security Suite — Zero-trust architecture rollout\n"
        "  3. Mobile Analytics Dashboard — iOS and Android release\n\n"
        "H1 2026 Targets:\n"
        "  1. AI Copilot for Enterprise — Natural language query interface\n"
        "  2. Edge Computing Framework — Low-latency processing nodes\n"
        "  3. Unified Customer Portal — Single sign-on across all products"
    )
    add_text_to_slide(slide5, 0.5, 1.3, 9, 5, roadmap_text, font_size=15)

    # ======= Slide 6: Key Clients =======
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_to_slide(slide6, 0.5, 0.3, 9, 0.8,
                      "Key Client Wins & Expansions",
                      font_size=28, bold=True, color=(0x1B, 0x3A, 0x6B))
    rows2, cols2 = 5, 3
    tbl2 = slide6.shapes.add_table(rows2, cols2,
                                    Inches(0.8), Inches(1.5),
                                    Inches(8.4), Inches(3))
    table2 = tbl2.table
    headers2 = ["Client", "Contract Value", "Status"]
    clients = [
        ["Pinnacle Healthcare", "$12.4M ARR", "New Win"],
        ["Eastwood Financial Group", "$8.7M ARR", "Expansion"],
        ["Nordic Logistics AB", "$6.2M ARR", "Renewal"],
        ["Cascade Energy Solutions", "$4.9M ARR", "New Win"],
    ]
    for c, h in enumerate(headers2):
        cell = table2.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row_data in enumerate(clients, 1):
        for c, val in enumerate(row_data):
            table2.cell(r, c).text = val

    # ======= Slide 7: Team & Hiring =======
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_to_slide(slide7, 0.5, 0.3, 9, 0.8,
                      "Team Growth & Talent Acquisition",
                      font_size=28, bold=True, color=(0x1B, 0x3A, 0x6B))
    team_text = (
        "Current Headcount: 847 employees (up from 762 in Q2)\n\n"
        "Key Hires:\n"
        "  - Dr. Anika Patel — Chief AI Officer (ex-DeepMind)\n"
        "  - Marcus Rivera — VP of Sales, EMEA Region\n"
        "  - Yuki Tanaka — Head of Cloud Infrastructure\n\n"
        "Open Positions: 43 roles across Engineering (28), Sales (9), "
        "and Operations (6). Average time-to-hire: 34 days, down from 41 in Q2.\n\n"
        "Employee Satisfaction Score: 4.3/5.0 (Glassdoor), "
        "voluntary attrition rate: 6.2% annualized."
    )
    add_text_to_slide(slide7, 0.5, 1.3, 9, 5, team_text, font_size=15)

    # ======= Slide 8: Financial Outlook =======
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_to_slide(slide8, 0.5, 0.3, 9, 0.8,
                      "Financial Outlook — FY 2025",
                      font_size=28, bold=True, color=(0x1B, 0x3A, 0x6B))
    finance_text = (
        "Full Year Revenue Forecast: $540M - $560M\n"
        "  - Cloud Services: $195M - $205M\n"
        "  - Enterprise Solutions: $160M - $165M\n"
        "  - Data Analytics: $108M - $115M\n"
        "  - Other Segments: $77M - $75M\n\n"
        "Projected Operating Margin: 21.5% - 23.0%\n"
        "Capital Expenditure Budget: $48M (data center expansion)\n"
        "R&D Investment: $72M (13.2% of revenue)\n"
        "Free Cash Flow Target: $85M - $95M"
    )
    add_text_to_slide(slide8, 0.5, 1.3, 9, 5, finance_text, font_size=15)

    # ======= Slide 9: Risk Assessment =======
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_to_slide(slide9, 0.5, 0.3, 9, 0.8,
                      "Risk Assessment & Mitigation",
                      font_size=28, bold=True, color=(0x1B, 0x3A, 0x6B))
    risk_text = (
        "High Priority Risks:\n"
        "  1. Supply chain disruption for hardware components — Mitigation: "
        "dual-supplier strategy with Foxconn and Pegatron\n"
        "  2. Regulatory compliance (EU AI Act) — Mitigation: dedicated "
        "compliance team, Q4 audit scheduled\n"
        "  3. Talent competition in AI/ML — Mitigation: enhanced RSU "
        "packages, remote work flexibility\n\n"
        "Medium Priority Risks:\n"
        "  1. Currency fluctuation (EUR/USD) — Natural hedge through "
        "European revenue streams\n"
        "  2. Customer concentration — Top 5 clients represent 28% of revenue, "
        "down from 35% last year"
    )
    add_text_to_slide(slide9, 0.5, 1.3, 9, 5, risk_text, font_size=14)

    # ======= Slide 10: Next Steps =======
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_to_slide(slide10, 0.5, 0.3, 9, 0.8,
                      "Next Steps & Action Items",
                      font_size=28, bold=True, color=(0x1B, 0x3A, 0x6B))
    steps_text = (
        "Immediate Actions (Next 30 Days):\n"
        "  - Finalize DataBridge v3.2 beta testing with 5 pilot customers\n"
        "  - Complete EMEA sales team onboarding (Marcus Rivera)\n"
        "  - Submit EU AI Act compliance preliminary report\n\n"
        "Q4 Milestones:\n"
        "  - Launch CloudGuard Security Suite (October 15)\n"
        "  - Close 3 additional enterprise deals ($15M+ pipeline)\n"
        "  - Complete data center expansion in Frankfurt\n\n"
        "Board Review: Scheduled for December 8, 2025\n"
        "Contact: strategy@meridiantech.com"
    )
    add_text_to_slide(slide10, 0.5, 1.3, 9, 5, steps_text, font_size=15)

    # Save presentation
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Create html_slides directory
    os.makedirs(HTML_DIR, exist_ok=True)
    print(f'Created directory: {HTML_DIR}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
