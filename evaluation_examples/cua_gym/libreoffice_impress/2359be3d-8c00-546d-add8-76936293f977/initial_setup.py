"""
Initial Setup: Department Review presentation with 9 slides, no master separator line.
Task ID: impress_ma_015
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_015'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide data for a realistic 9-slide department review
    slides_data = [
        {
            "layout": 0,  # Title Slide
            "title": "Q1 2025 Department Review",
            "subtitle": "Prepared by the Operations Team\nMarch 28, 2025",
        },
        {
            "layout": 1,  # Title + Content
            "title": "Executive Summary",
            "content": (
                "Overall revenue grew 12% year-over-year to $4.8M\n"
                "Customer acquisition cost decreased by 8%\n"
                "Net Promoter Score improved from 42 to 51\n"
                "Three new product features launched on schedule\n"
                "Headcount increased by 6 across Engineering and Sales"
            ),
        },
        {
            "layout": 1,
            "title": "Engineering Department",
            "content": (
                "Sprint velocity improved 15% quarter-over-quarter\n"
                "Deployed 47 production releases with zero P0 incidents\n"
                "Reduced average build time from 18 min to 9 min\n"
                "Completed migration of legacy auth service to microservices\n"
                "Hired 3 senior engineers (backend, ML, platform)"
            ),
        },
        {
            "layout": 1,
            "title": "Sales & Marketing",
            "content": (
                "Closed 23 new enterprise deals worth $1.9M total\n"
                "Marketing-qualified leads increased 28% via content strategy\n"
                "Average deal cycle shortened from 45 to 38 days\n"
                "Launched partner referral program with 12 initial partners\n"
                "Social media engagement up 34% across all channels"
            ),
        },
        {
            "layout": 1,
            "title": "Customer Success",
            "content": (
                "Churn rate decreased from 4.2% to 3.1%\n"
                "Average ticket resolution time reduced to 4.2 hours\n"
                "Onboarding completion rate improved to 91%\n"
                "Launched self-service knowledge base with 150+ articles\n"
                "Upsell revenue contributed $620K to quarterly total"
            ),
        },
        {
            "layout": 1,
            "title": "Finance & Operations",
            "content": (
                "Operating margin improved to 18.5% from 15.2%\n"
                "Accounts receivable days reduced from 52 to 41\n"
                "Implemented new expense approval workflow saving 12 hrs/week\n"
                "Completed SOC 2 Type II audit with zero findings\n"
                "Renegotiated cloud hosting contract saving $180K annually"
            ),
        },
        {
            "layout": 1,
            "title": "People & Culture",
            "content": (
                "Employee satisfaction score: 4.3/5.0 (up from 3.9)\n"
                "Voluntary turnover rate: 8% (industry avg: 13%)\n"
                "Launched mentorship program with 28 pairs\n"
                "Completed diversity & inclusion training for all managers\n"
                "Remote work policy updated based on employee survey feedback"
            ),
        },
        {
            "layout": 1,
            "title": "Key Risks & Challenges",
            "content": (
                "Supply chain delays affecting hardware procurement timelines\n"
                "Increasing competition in the mid-market segment\n"
                "Technical debt in payment processing module needs attention\n"
                "Two senior leadership positions remain unfilled\n"
                "Regulatory changes in EU market require compliance updates by Q3"
            ),
        },
        {
            "layout": 1,
            "title": "Q2 2025 Priorities",
            "content": (
                "Launch mobile app v2.0 with offline capabilities\n"
                "Expand into APAC market with localized offerings\n"
                "Achieve $5.5M revenue target (+15% QoQ)\n"
                "Complete data platform migration to new infrastructure\n"
                "Fill remaining leadership roles by end of April"
            ),
        },
    ]

    for i, sd in enumerate(slides_data):
        layout_idx = sd["layout"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        if slide.shapes.title:
            slide.shapes.title.text = sd["title"]

        if layout_idx == 0 and "subtitle" in sd:
            # Title slide: subtitle in placeholder 1
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = sd["subtitle"]
        elif "content" in sd:
            # Content slides: body in placeholder 1
            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                lines = sd["content"].split("\n")
                for j, line in enumerate(lines):
                    if j == 0:
                        tf.paragraphs[0].text = line
                    else:
                        p = tf.add_paragraph()
                        p.text = line
                        p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
