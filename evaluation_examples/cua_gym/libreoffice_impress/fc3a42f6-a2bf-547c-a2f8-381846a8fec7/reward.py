"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 41, I want to pop in a 5-point star that’s exactly 2 cm by 2 cm, tucked right up in the top-left corner. Fill it with the Purple 3 swatch (hex #800080). How do I set that up in LibreOffice Impress?
Generated: 2025-09-10 22:43:32
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
import glob
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE

# ------------------------------------------------------------
# Reward Script: Verify a 2 cm x 2 cm, Purple (#800080) 5-point
# star positioned at the top-left corner of slide 41.
# ------------------------------------------------------------
# Scoring rubric (progressive):
#   +0.3  – Star exists on slide 41
#   +0.2  – Star size is 2 cm × 2 cm (±0.1 cm tolerance)
#   +0.2  – Star positioned within 0.1 cm of (0,0)
#   +0.3  – Star fill colour exactly #800080
# Total possible: 1.0
# ------------------------------------------------------------

EMU_PER_CM = 360000  # English Metric Units per centimetre (pptx internal)


def _cm(emu):
    """Convert EMU to centimetres (float)."""
    return emu / EMU_PER_CM


def _find_presentation():
    """Return path of the first .pptx file in /home/user or its sub-dirs."""
    explicit = "/home/user/on_slide_41_i_want_to_pop_in_a_5_point_star_thats_exactly_2_cm_by_2_cm_tucked_right_up_in_the_top_le_golden.pptx"
    if os.path.exists(explicit):
        return explicit
    matches = glob.glob("/home/user/**/*.pptx", recursive=True)
    return matches[0] if matches else None


def verify_star_on_slide_41(file_path: str) -> float:
    """Verify task completion and return a reward score between 0.0 and 1.0."""
    MAX_SCORE = 1.0
    score = 0.0

    # --- Load presentation -------------------------------------------------
    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded presentation: {file_path}  (slides: {len(prs.slides)})")
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # --- Ensure slide 41 exists -------------------------------------------
    target_index = 40  # 0-based index for slide 41
    if len(prs.slides) <= target_index:
        print(f"✗ Slide 41 missing (found only {len(prs.slides)} slide(s))")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[target_index]

    # --- Locate 5-point star(s) -------------------------------------------
    stars = [sh for sh in slide.shapes
             if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and
             sh.auto_shape_type == MSO_AUTO_SHAPE_TYPE.STAR_5_POINT]

    if not stars:
        print("✗ No 5-point star found on slide 41")
        print("REWARD: 0.0")
        return 0.0

    # Iterate through stars and keep the best-scoring one
    best_local = 0.0
    for sh in stars:
        local = 0.3  # existence points
        width_cm, height_cm = _cm(sh.width), _cm(sh.height)
        left_cm, top_cm = _cm(sh.left), _cm(sh.top)

        # Size check --------------------------------------------------------
        if abs(width_cm - 2.0) <= 0.1 and abs(height_cm - 2.0) <= 0.1:
            local += 0.2
            size_note = "✓ Size within tolerance"
        else:
            size_note = "✗ Size out of tolerance"

        # Position check ----------------------------------------------------
        if left_cm <= 0.1 and top_cm <= 0.1:
            local += 0.2
            pos_note = "✓ Position correct (top-left)"
        else:
            pos_note = "✗ Position incorrect"

        # Colour check ------------------------------------------------------
        color_ok = False
        try:
            fill = sh.fill
            if fill and fill.fore_color.type == 1:  # RGB colour
                rgb = fill.fore_color.rgb  # returns (r, g, b)
                if rgb and (rgb[0], rgb[1], rgb[2]) == (128, 0, 128):
                    color_ok = True
        except Exception:
            pass

        if color_ok:
            local += 0.3
            colour_note = "✓ Colour matches #800080"
        else:
            colour_note = "✗ Colour does not match #800080"

        # Debug summary for this star
        print("--- Star candidate ---")
        print(f"  Size       : {width_cm:.2f} cm × {height_cm:.2f} cm   {size_note}")
        print(f"  Position   : left {left_cm:.2f} cm, top {top_cm:.2f} cm   {pos_note}")
        print(f"  Colour     : {colour_note}")
        print(f"  Partial score: {local}\n")

        best_local = max(best_local, local)

    score = min(best_local, MAX_SCORE)
    print(f"Final score: {score}/{MAX_SCORE}")
    print(f"REWARD: {score}")
    return score


if __name__ == "__main__":
    pptx_path = _find_presentation()
    if pptx_path:
        verify_star_on_slide_41(pptx_path)
    else:
        print("✗ No .pptx file found to evaluate")
        print("REWARD: 0.0")
