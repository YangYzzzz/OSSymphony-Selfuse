"""
Initial Setup: Resize and reposition bar chart on slide 3
Task ID: impress_gf3_042
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_042'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard widescreen: 33.87 cm x 19.05 cm
    prs.slide_width = Cm(33.87)
    prs.slide_height = Cm(19.05)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q3 2025 Sales Performance Review"
    slide1.placeholders[1].text = "Regional Analysis & Growth Metrics\nPrepared by Analytics Team"

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    body2 = slide2.placeholders[1]
    tf2 = body2.text_frame
    tf2.text = "Key highlights from the third quarter:"
    bullets = [
        "Total revenue increased 12.4% year-over-year to $4.87M",
        "Western region led growth with 18.2% increase",
        "New product line contributed $620K in first quarter of availability",
        "Customer acquisition cost decreased by 7.3%",
        "Net promoter score improved from 72 to 78",
    ]
    for b in bullets:
        p = tf2.add_paragraph()
        p.text = b
        p.level = 1

    # --- Slide 3: Bar Chart (initial state: ~14cm wide, ~9cm tall, offset) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title textbox
    txBox = slide3.shapes.add_textbox(Cm(2), Cm(0.5), Cm(28), Cm(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Regional Revenue Comparison"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x3A, 0x4F)

    # Bar chart data
    chart_data = CategoryChartData()
    chart_data.categories = ['North', 'South', 'East', 'West', 'Central']
    chart_data.add_series('Q2 2025', (980, 745, 1120, 1340, 560))
    chart_data.add_series('Q3 2025', (1050, 820, 1180, 1585, 635))

    # Place chart at offset position: ~14cm wide, ~9cm tall
    # Offset from center: left at ~4cm, top at ~6cm
    chart_frame = slide3.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Cm(4.0), Cm(6.0),   # left, top (offset, not centered)
        Cm(14.0), Cm(9.0),  # width, height
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.include_in_layout = False

    # --- Slide 4: Table Data ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox4 = slide4.shapes.add_textbox(Cm(2), Cm(0.5), Cm(28), Cm(2))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Detailed Revenue Breakdown"
    p4.alignment = PP_ALIGN.CENTER
    run4 = p4.runs[0]
    run4.font.size = Pt(22)
    run4.font.bold = True

    table_shape = slide4.shapes.add_table(
        6, 4, Cm(3), Cm(3), Cm(27), Cm(12)
    )
    table = table_shape.table
    headers = ['Region', 'Q2 Revenue ($K)', 'Q3 Revenue ($K)', 'Growth (%)']
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)

    rows_data = [
        ['North', '980', '1,050', '7.1%'],
        ['South', '745', '820', '10.1%'],
        ['East', '1,120', '1,180', '5.4%'],
        ['West', '1,340', '1,585', '18.3%'],
        ['Central', '560', '635', '13.4%'],
    ]
    for r, row_data in enumerate(rows_data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 5: Key Insights ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Key Insights"
    body5 = slide5.placeholders[1]
    tf5 = body5.text_frame
    tf5.text = "Growth Drivers"
    insights = [
        "Western region expansion fueled by new distribution partnerships",
        "Southern region recovery after Q1 supply chain disruptions",
        "Central region benefiting from recent marketing campaign",
        "Eastern region stable but approaching market saturation",
    ]
    for ins in insights:
        p = tf5.add_paragraph()
        p.text = ins
        p.level = 1

    # --- Slide 6: Next Steps ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Q4 Action Items"
    body6 = slide6.placeholders[1]
    tf6 = body6.text_frame
    tf6.text = "Priority initiatives for Q4 2025:"
    actions = [
        "Expand western region warehouse capacity by November",
        "Launch targeted campaign for eastern region market penetration",
        "Onboard 3 new enterprise clients in the central region",
        "Implement dynamic pricing model across all regions",
        "Complete CRM migration by end of October",
    ]
    for a in actions:
        p = tf6.add_paragraph()
        p.text = a
        p.level = 1

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
