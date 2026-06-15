"""
Reward Script: Resize and reposition bar chart on slide 3
Task ID: impress_gf3_042
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30) - Chart width is 18 cm
  Component 2 (0.30) - Chart height is 11 cm
  Component 3 (0.20) - Chart top edge at 5.5 cm
  Component 4 (0.20) - Chart horizontally centered on slide
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_042'

# Tolerance: 0.1 cm = 36000 EMU
TOLERANCE_EMU = 36000

# Target values in EMU (1 cm = 360000 EMU)
TARGET_WIDTH = 6480000    # 18 cm
TARGET_HEIGHT = 3960000   # 11 cm
TARGET_TOP = 1980000      # 5.5 cm


def persist_app_state(domain):
    """Save any unsaved LibreOffice edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]

    # Find the chart shape on slide 3
    chart_shape = None
    for shape in slide3.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            chart_shape = shape
            break

    if chart_shape is None:
        print("FAIL: No chart found on slide 3")
        print("REWARD: 0.0")
        return 0.0

    slide_width = prs.slide_width
    print(f"INFO: Slide width = {slide_width} EMU ({round(slide_width/360000, 2)} cm)")
    print(f"INFO: Chart — left={chart_shape.left}, top={chart_shape.top}, width={chart_shape.width}, height={chart_shape.height}")
    print(f"INFO: Chart — left={round(chart_shape.left/360000,2)} cm, top={round(chart_shape.top/360000,2)} cm, "
          f"width={round(chart_shape.width/360000,2)} cm, height={round(chart_shape.height/360000,2)} cm")

    # Component 1: Chart width is 18 cm (0.30 points)
    try:
        width_diff = abs(chart_shape.width - TARGET_WIDTH)
        if width_diff <= TOLERANCE_EMU:
            print(f"PASS: Component 1 — Chart width is {round(chart_shape.width/360000,2)} cm "
                  f"(target 18.0 cm, diff {round(width_diff/360000,3)} cm) (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 1 — Chart width is {round(chart_shape.width/360000,2)} cm, "
                  f"expected 18.0 cm (diff {round(width_diff/360000,3)} cm, tolerance 0.1 cm)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Chart height is 11 cm (0.30 points)
    try:
        height_diff = abs(chart_shape.height - TARGET_HEIGHT)
        if height_diff <= TOLERANCE_EMU:
            print(f"PASS: Component 2 — Chart height is {round(chart_shape.height/360000,2)} cm "
                  f"(target 11.0 cm, diff {round(height_diff/360000,3)} cm) (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 2 — Chart height is {round(chart_shape.height/360000,2)} cm, "
                  f"expected 11.0 cm (diff {round(height_diff/360000,3)} cm, tolerance 0.1 cm)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Chart top edge at 5.5 cm (0.20 points)
    try:
        top_diff = abs(chart_shape.top - TARGET_TOP)
        if top_diff <= TOLERANCE_EMU:
            print(f"PASS: Component 3 — Chart top is {round(chart_shape.top/360000,2)} cm "
                  f"(target 5.5 cm, diff {round(top_diff/360000,3)} cm) (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 3 — Chart top is {round(chart_shape.top/360000,2)} cm, "
                  f"expected 5.5 cm (diff {round(top_diff/360000,3)} cm, tolerance 0.1 cm)")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Chart horizontally centered on slide (0.20 points)
    # Expected left = (slide_width - chart_width) / 2
    try:
        expected_left = (slide_width - chart_shape.width) // 2
        left_diff = abs(chart_shape.left - expected_left)
        # Use tolerance of 0.1 cm for centering as well
        if left_diff <= TOLERANCE_EMU:
            print(f"PASS: Component 4 — Chart is horizontally centered "
                  f"(left={round(chart_shape.left/360000,2)} cm, expected={round(expected_left/360000,2)} cm, "
                  f"diff {round(left_diff/360000,3)} cm) (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 4 — Chart not centered "
                  f"(left={round(chart_shape.left/360000,2)} cm, expected={round(expected_left/360000,2)} cm, "
                  f"diff {round(left_diff/360000,3)} cm, tolerance 0.1 cm)")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
