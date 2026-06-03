"""
Reward Script: Apply strikethrough formatting AND change text color to gray (#808080)
                on the first 2 bullet points of the action items list on slide 4.
Task ID: osworld_impress_strikethrough_text_007
Domain: libreoffice_impress
Scoring:
  Component 1: Bullet 1 (para 0) has strikethrough (sngStrike)        — 0.30 pts
  Component 2: Bullet 1 (para 0) text color is gray (#808080)         — 0.20 pts
  Component 3: Bullet 2 (para 1) has strikethrough (sngStrike)        — 0.30 pts
  Component 4: Bullet 2 (para 1) text color is gray (#808080)         — 0.20 pts
  Total: 1.0
"""

import os

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_strikethrough_text_007'

EXPECTED_GRAY = '808080'
SLIDE_IDX = 3       # slide 4 is 0-indexed as 3
BODY_SHAPE_NAME = 'Content Placeholder 2'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0

    The task: apply strikethrough formatting AND change text color to gray (#808080)
    on the first 2 bullet points of the action items list on slide 4.

    Only paragraphs 0 and 1 of the body textbox on slide 4 should be modified.
    """
    total_score = 0.0

    # Precondition: load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"CRITICAL: Expected at least 4 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[SLIDE_IDX]

    # Find the body text shape containing bullet points
    body_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == BODY_SHAPE_NAME:
            body_shape = shape
            break

    if body_shape is None:
        # Fallback: find any shape that is NOT the title and has multiple paragraphs
        for shape in slide.shapes:
            if shape.has_text_frame and len(shape.text_frame.paragraphs) >= 2:
                title_shape = slide.shapes.title
                if title_shape is None or shape != title_shape:
                    body_shape = shape
                    break

    if body_shape is None:
        print("CRITICAL: Cannot find body text shape on slide 4")
        print("REWARD: 0.0")
        return 0.0

    paragraphs = body_shape.text_frame.paragraphs

    if len(paragraphs) < 2:
        print(f"CRITICAL: Expected at least 2 bullet points, found {len(paragraphs)}")
        print("REWARD: 0.0")
        return 0.0

    # Helper: get effective runs (non-empty text) for a paragraph
    def get_runs(para):
        return [r for r in para.runs if (r.text or "").strip()]

    # Helper: check strikethrough on a paragraph's runs
    def has_strikethrough(para):
        runs = get_runs(para)
        if not runs:
            return False
        # All non-empty runs must have sngStrike
        return all(r.font._element.attrib.get('strike', 'noStrike') == 'sngStrike' for r in runs)

    # Helper: check gray color on a paragraph's runs
    def has_gray_color(para):
        runs = get_runs(para)
        if not runs:
            return False
        # All non-empty runs must have color type set and be #808080
        for r in runs:
            try:
                if r.font.color.type is None:
                    return False
                rgb = str(r.font.color.rgb).upper()
                if rgb != EXPECTED_GRAY.upper():
                    return False
            except Exception:
                return False
        return True

    para0 = paragraphs[0]
    para1 = paragraphs[1]

    # Component 1: Bullet 1 (para 0) has strikethrough (0.30 pts)
    try:
        if has_strikethrough(para0):
            print(f"PASS: Component 1 — Bullet 1 has strikethrough (sngStrike) (0.30 pts)")
            total_score += 0.30
        else:
            runs = get_runs(para0)
            strike_vals = [r.font._element.attrib.get('strike', 'noStrike') for r in runs]
            print(f"FAIL: Component 1 — Bullet 1 strikethrough expected sngStrike, found {strike_vals}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Bullet 1 (para 0) text color is gray #808080 (0.20 pts)
    try:
        if has_gray_color(para0):
            print(f"PASS: Component 2 — Bullet 1 color is gray (#808080) (0.20 pts)")
            total_score += 0.20
        else:
            runs = get_runs(para0)
            colors = []
            for r in runs:
                try:
                    colors.append(str(r.font.color.rgb) if r.font.color.type is not None else 'None')
                except Exception:
                    colors.append('ERROR')
            print(f"FAIL: Component 2 — Bullet 1 color expected #808080, found {colors}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Bullet 2 (para 1) has strikethrough (0.30 pts)
    try:
        if has_strikethrough(para1):
            print(f"PASS: Component 3 — Bullet 2 has strikethrough (sngStrike) (0.30 pts)")
            total_score += 0.30
        else:
            runs = get_runs(para1)
            strike_vals = [r.font._element.attrib.get('strike', 'noStrike') for r in runs]
            print(f"FAIL: Component 3 — Bullet 2 strikethrough expected sngStrike, found {strike_vals}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Bullet 2 (para 1) text color is gray #808080 (0.20 pts)
    try:
        if has_gray_color(para1):
            print(f"PASS: Component 4 — Bullet 2 color is gray (#808080) (0.20 pts)")
            total_score += 0.20
        else:
            runs = get_runs(para1)
            colors = []
            for r in runs:
                try:
                    colors.append(str(r.font.color.rgb) if r.font.color.type is not None else 'None')
                except Exception:
                    colors.append('ERROR')
            print(f"FAIL: Component 4 — Bullet 2 color expected #808080, found {colors}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
