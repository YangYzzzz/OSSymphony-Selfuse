"""
Initial Setup: 6-slide quarterly review presentation — pre-task state
Task ID: osworld_impress_strikethrough_text_007
Domain: libreoffice_impress

Slide 4 has 5 action-item bullet points in plain black text (no strikethrough,
no gray color) so the agent can apply the formatting in the task.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_strikethrough_text_007'
OUTPUT  = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Use standard 16:9 widescreen dimensions
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    layouts = prs.slide_layouts
    # 0 = Title Slide, 1 = Title+Content, 5 = Blank, 6 = Title Only

    # ── Slide 1: Title Slide ──────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(layouts[0])
    slide1.shapes.title.text = "Q2 2025 Quarterly Review"
    slide1.placeholders[1].text = "Strategic Planning & Performance Overview\nApril – June 2025"

    # ── Slide 2: Executive Summary ────────────────────────────────────────────
    slide2 = prs.slides.add_slide(layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Revenue grew 12% YoY to $4.7 M in Q2"
    bullets2 = [
        "Customer acquisition cost reduced by 8% through targeted campaigns",
        "Product launch of ProSuite v3.0 exceeded adoption targets by 22%",
        "Headcount expanded to 148 FTEs across 4 regional offices",
        "Net Promoter Score improved from 42 to 57 quarter-over-quarter",
    ]
    for b in bullets2:
        p = tf2.add_paragraph()
        p.text = b
        p.level = 1

    # ── Slide 3: Financial Highlights ────────────────────────────────────────
    slide3 = prs.slides.add_slide(layouts[1])
    slide3.shapes.title.text = "Financial Highlights"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Key Metrics — Q2 2025 vs Q2 2024"
    fin_items = [
        "Total Revenue: $4.7 M  (+12% YoY)",
        "Gross Margin: 63%  (+3 pp YoY)",
        "Operating Expenses: $2.1 M  (-5% YoY)",
        "EBITDA: $890 K  (+28% YoY)",
        "Cash & Equivalents: $3.2 M (healthy runway)",
    ]
    for item in fin_items:
        p = tf3.add_paragraph()
        p.text = item
        p.level = 1

    # ── Slide 4: Action Items (THE KEY SLIDE) ────────────────────────────────
    # All 5 bullet points must be in plain black text with NO strikethrough.
    slide4 = prs.slides.add_slide(layouts[1])
    slide4.shapes.title.text = "Action Items"

    tf4 = slide4.placeholders[1].text_frame
    tf4.clear()  # remove any default placeholder text

    action_items = [
        "Schedule follow-up meetings with top 10 enterprise accounts by July 15",
        "Submit revised marketing budget proposal to CFO before end of month",
        "Complete onboarding for 6 new engineers joining the platform team",
        "Finalize partnership agreement with DataBridge Technologies by August 1",
        "Deploy ProSuite v3.1 hotfix to production environment this sprint",
    ]

    for idx, item_text in enumerate(action_items):
        if idx == 0:
            p = tf4.paragraphs[0]
        else:
            p = tf4.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = item_text
        # Explicitly set black color and NO strikethrough for all items
        run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        run.font.size = Pt(18)
        run.font.bold = False
        # Ensure strike attribute is absent / set to noStrike
        run.font._element.attrib['strike'] = 'noStrike'

    # ── Slide 5: Risks & Mitigations ─────────────────────────────────────────
    slide5 = prs.slides.add_slide(layouts[1])
    slide5.shapes.title.text = "Risks & Mitigations"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Identified Risks for Q3 2025"
    risks = [
        "Supply chain delays: engage 2nd-tier supplier network proactively",
        "Regulatory changes in EU market: legal review scheduled for July 20",
        "Key talent retention: comp review and equity refresh planned for August",
        "Competitor product launch: accelerate feature roadmap Q3 milestones",
    ]
    for r in risks:
        p = tf5.add_paragraph()
        p.text = r
        p.level = 1

    # ── Slide 6: Next Steps & Closing ─────────────────────────────────────────
    slide6 = prs.slides.add_slide(layouts[0])
    slide6.shapes.title.text = "Next Steps"
    slide6.placeholders[1].text = (
        "Q3 Planning Kick-off: July 8, 2025\n"
        "All-Hands Town Hall: July 22, 2025\n"
        "Mid-Quarter Review: August 19, 2025"
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup — open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
