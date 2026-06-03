"""
Reward Script: Copy the title of slide 2 into the notes section of slide 2.
Task ID: osworld_impress_slide_notes_003
Domain: libreoffice_impress
Scoring:
  Component 1 (0.6): Slide 2 notes area contains the text 'Market Analysis'
  Component 2 (0.4): Slide 2 notes area matches exactly 'Market Analysis' (no extra text)
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_slide_notes_003'


def get_slide_notes(slide):
    """Return stripped notes text for a slide, or empty string on error."""
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def get_slide_title(slide):
    """Return the title placeholder text for a slide, or empty string."""
    try:
        for shape in slide.shapes:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                if shape.placeholder_format.idx == 0:
                    return shape.text_frame.text.strip()
    except Exception:
        pass
    return ""


def verify_task(file_path):
    """
    Verify that slide 2's notes area contains its title text 'Market Analysis'.
    Returns a float between 0.0 and 1.0.
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: ensure we have at least 2 slides
    if len(prs.slides) < 2:
        print(f"CRITICAL: Expected at least 2 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]

    # Get slide 2 title for reference
    slide2_title = get_slide_title(slide2)
    print(f"INFO: Slide 2 title: {repr(slide2_title)}")

    # Get slide 2 notes
    slide2_notes = get_slide_notes(slide2)
    print(f"INFO: Slide 2 notes: {repr(slide2_notes)}")

    # Expected text (from task context: slide 2 title is 'Market Analysis')
    expected_text = "Market Analysis"

    # Component 1: Slide 2 notes area contains 'Market Analysis' (0.6 points)
    # This FAILS on initial_env (notes is empty) and PASSES on golden_env
    try:
        if expected_text in slide2_notes:
            print(f"PASS: Component 1 — Slide 2 notes contains '{expected_text}' (0.6 pts)")
            total_score += 0.6
        else:
            print(f"FAIL: Component 1 — Slide 2 notes does not contain '{expected_text}', "
                  f"found: {repr(slide2_notes)}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 2 notes matches exactly 'Market Analysis' (0.4 points)
    # Awards points only when the copy is exact (no extra text added/appended)
    # This FAILS on initial_env (notes is empty) and PASSES on golden_env
    try:
        if slide2_notes == expected_text:
            print(f"PASS: Component 2 — Slide 2 notes exactly matches '{expected_text}' (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 2 — Slide 2 notes is not exactly '{expected_text}', "
                  f"found: {repr(slide2_notes)}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
