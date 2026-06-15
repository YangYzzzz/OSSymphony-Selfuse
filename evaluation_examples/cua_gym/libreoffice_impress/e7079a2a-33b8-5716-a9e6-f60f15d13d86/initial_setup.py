"""
Initial Setup: Copy the title of slide 2 into the notes section of slide 2.
Task ID: osworld_impress_slide_notes_003
Domain: libreoffice_impress

Creates a 5-slide business overview deck. Slide 2 has the title 'Market Analysis'.
The notes area of slide 2 is intentionally empty (the agent must fill it).
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_slide_notes_003'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # Default slide dimensions (10 x 7.5 inches landscape)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    layout_title = prs.slide_layouts[0]  # Title Slide layout
    slide1 = prs.slides.add_slide(layout_title)
    slide1.shapes.title.text = "Nexus Solutions Inc."
    slide1.placeholders[1].text = "Annual Business Overview 2025"

    # ---- Slide 2: Market Analysis (title only, notes empty) ----
    layout_tc = prs.slide_layouts[1]  # Title and Content layout
    slide2 = prs.slides.add_slide(layout_tc)
    slide2.shapes.title.text = "Market Analysis"
    content2 = slide2.placeholders[1]
    tf2 = content2.text_frame
    tf2.text = "Total Addressable Market: $4.2B"
    p2 = tf2.add_paragraph()
    p2.text = "Year-over-Year Growth: 18.7%"
    p2.level = 0
    p3 = tf2.add_paragraph()
    p3.text = "Primary Segments: Enterprise (62%), SMB (28%), Consumer (10%)"
    p3.level = 0
    p4 = tf2.add_paragraph()
    p4.text = "Key Competitors: 3 major players with combined 45% market share"
    p4.level = 0
    # Notes for slide 2 must be EMPTY — task is to add the title there
    # Do NOT access slide2.notes_slide here (accessing it creates a notes placeholder)

    # ---- Slide 3: Revenue Performance ----
    slide3 = prs.slides.add_slide(layout_tc)
    slide3.shapes.title.text = "Revenue Performance"
    content3 = slide3.placeholders[1]
    tf3 = content3.text_frame
    tf3.text = "Q1 Revenue: $12.4M (+22% YoY)"
    p3b = tf3.add_paragraph()
    p3b.text = "Q2 Revenue: $15.1M (+19% YoY)"
    p3c = tf3.add_paragraph()
    p3c.text = "Q3 Revenue: $17.8M (+25% YoY)"
    p3d = tf3.add_paragraph()
    p3d.text = "Q4 Forecast: $21.3M (+28% YoY)"
    slide3.notes_slide.notes_text_frame.text = "Discuss Q3 outperformance drivers with board."

    # ---- Slide 4: Strategic Initiatives ----
    slide4 = prs.slides.add_slide(layout_tc)
    slide4.shapes.title.text = "Strategic Initiatives"
    content4 = slide4.placeholders[1]
    tf4 = content4.text_frame
    tf4.text = "Initiative 1: APAC Market Expansion — Target 15 new enterprise accounts"
    p4b = tf4.add_paragraph()
    p4b.text = "Initiative 2: Product Platform Unification — Reduce SKU count by 40%"
    p4c = tf4.add_paragraph()
    p4c.text = "Initiative 3: Talent Acquisition — Hire 85 engineers by Q3"
    p4d = tf4.add_paragraph()
    p4d.text = "Initiative 4: Partner Ecosystem — Onboard 20 certified resellers"
    slide4.notes_slide.notes_text_frame.text = "Focus on APAC timeline milestones."

    # ---- Slide 5: Conclusion & Next Steps ----
    slide5 = prs.slides.add_slide(layout_tc)
    slide5.shapes.title.text = "Conclusion & Next Steps"
    content5 = slide5.placeholders[1]
    tf5 = content5.text_frame
    tf5.text = "Strong market position with accelerating revenue trajectory"
    p5b = tf5.add_paragraph()
    p5b.text = "Execute on 4 strategic initiatives through H2 2025"
    p5c = tf5.add_paragraph()
    p5c.text = "Next Board Review: September 15, 2025"
    p5d = tf5.add_paragraph()
    p5d.text = "Contact: strategy@nexussolutions.com"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup — open in LibreOffice Impress with DISPLAY=:0
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
