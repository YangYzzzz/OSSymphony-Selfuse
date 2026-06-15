"""
Initial Setup: Award ceremony presentation with 8 slides, no transitions
Task ID: impress_tm_012
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_012'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def add_blank_slide_with_title(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    slide1 = add_title_slide(
        prs,
        "Annual Excellence Awards 2025",
        "Celebrating Outstanding Achievements Across All Departments"
    )

    # Slide 2: Agenda
    slide2 = add_content_slide(prs, "Ceremony Agenda", [
        "Welcome and Opening Remarks by CEO Patricia Huang",
        "Review of 2025 Achievements and Milestones",
        "Category Award Announcements",
        "Special Recognition: Lifetime Achievement Award",
        "Closing Remarks and Reception Details",
    ])

    # Slide 3: "And the Winner Is..." - NO transition, NO sound
    slide3 = add_content_slide(prs, "And the Winner Is...", [
        "Best Innovation Award",
        "Nominees:",
        "  - Project Aurora (Engineering Team)",
        "  - Customer Insights Platform (Analytics Division)",
        "  - GreenTech Initiative (Sustainability Group)",
        "",
        "The envelope, please...",
    ])

    # Slide 4: Winner Announcement
    slide4 = add_content_slide(prs, "Project Aurora - Engineering Team", [
        "Led by Dr. Rajesh Patel and team of 12 engineers",
        "Developed breakthrough AI-powered quality assurance system",
        "Reduced defect rates by 47% in Q3 alone",
        "Saved an estimated $2.3 million in production costs",
        "Patent pending on core algorithm",
    ])

    # Slide 5: Lifetime Achievement
    slide5 = add_content_slide(prs, "Lifetime Achievement Award", [
        "Presented to: Margaret O'Sullivan",
        "28 years of dedicated service",
        "Pioneered the company's international expansion strategy",
        "Mentored over 150 junior employees",
        "Founded the Women in Tech internal program",
    ])

    # Slide 6: Team Awards
    slide6 = add_content_slide(prs, "Team Excellence Awards", [
        "Customer Success Team - 98.7% satisfaction rating",
        "DevOps Engineering - 99.99% uptime achievement",
        "Marketing Creative - Campaign of the Year",
        "Finance Operations - Zero-error quarterly close",
    ])

    # Slide 7: Looking Ahead
    slide7 = add_content_slide(prs, "Looking Ahead to 2026", [
        "New innovation lab opening in Austin, TX",
        "Expanded mentorship programs across all regions",
        "Sustainability targets: carbon neutral by 2027",
        "Employee wellness initiative launch in Q1",
        "Global hackathon series: 4 events across 3 continents",
    ])

    # Slide 8: Closing
    slide8 = add_title_slide(
        prs,
        "Thank You & Congratulations!",
        "Join us for the reception in the Grand Ballroom\nRefreshments served from 7:00 PM"
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
