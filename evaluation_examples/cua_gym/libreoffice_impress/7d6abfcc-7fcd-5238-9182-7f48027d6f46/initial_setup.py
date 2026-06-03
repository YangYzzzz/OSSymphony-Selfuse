"""
Initial Setup: 7-slide product strategy deck — slide 5 title is 'Competition' (left-aligned)
Task ID: osworld_impress_title_set_aligned_005
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_set_aligned_005'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_content_slide(prs, title_text, body_lines, title_align=PP_ALIGN.LEFT):
    """Add a slide with title and bullet content."""
    slide_layout = prs.slide_layouts[1]  # Title + Content layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title_text
    # Apply alignment to all paragraphs in title
    for para in title_shape.text_frame.paragraphs:
        para.alignment = title_align

    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = body_lines[0]
    for line in body_lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0

    return slide


def create_initial():
    prs = Presentation()

    # ---- Slide 1: Title Slide ----
    slide1_layout = prs.slide_layouts[0]  # Title Slide layout
    slide1 = prs.slides.add_slide(slide1_layout)
    slide1.shapes.title.text = "NovaTech X1 Product Strategy"
    slide1.placeholders[1].text = "Q3 2025 — Roadmap & Market Positioning"

    # ---- Slide 2: Executive Summary ----
    add_title_content_slide(
        prs,
        title_text="Executive Summary",
        body_lines=[
            "NovaTech X1 targets the mid-market enterprise segment",
            "Projected 18% revenue growth in 12 months",
            "Three core pillars: Performance, Reliability, Integration",
            "TAM estimated at $4.2B globally",
        ],
    )

    # ---- Slide 3: Market Opportunity ----
    add_title_content_slide(
        prs,
        title_text="Market Opportunity",
        body_lines=[
            "Segment growth rate: 22% YoY (2024–2027)",
            "Key verticals: Healthcare, Finance, Manufacturing",
            "Customer pain points: legacy system costs, poor UX",
            "Underserved SMB tier presents $800M+ greenfield",
        ],
    )

    # ---- Slide 4: Product Roadmap ----
    add_title_content_slide(
        prs,
        title_text="Product Roadmap",
        body_lines=[
            "Q3 2025 — Beta launch with 50 design-partner customers",
            "Q4 2025 — GA release with full API suite",
            "Q1 2026 — Mobile companion app (iOS/Android)",
            "Q2 2026 — AI-powered analytics module",
        ],
    )

    # ---- Slide 5: Competition (left-aligned title — task target) ----
    slide5_layout = prs.slide_layouts[1]
    slide5 = prs.slides.add_slide(slide5_layout)

    title5 = slide5.shapes.title
    title5.text = "Competition"
    # Explicitly set left alignment for title paragraph
    for para in title5.text_frame.paragraphs:
        para.alignment = PP_ALIGN.LEFT

    body5 = slide5.placeholders[1]
    tf5 = body5.text_frame
    tf5.text = "Direct competitors: Apex Systems, CoreEdge, Prism Software"
    bullets5 = [
        "Apex Systems: strong brand, high price point, limited integrations",
        "CoreEdge: broad feature set, but poor customer support ratings",
        "Prism Software: niche player, healthcare-focused, weak roadmap",
        "NovaTech X1 differentiator: open API, transparent pricing, 24/7 SLA",
    ]
    for line in bullets5:
        p = tf5.add_paragraph()
        p.text = line
        p.level = 0

    # ---- Slide 6: Go-to-Market Strategy ----
    add_title_content_slide(
        prs,
        title_text="Go-to-Market Strategy",
        body_lines=[
            "Channel: Direct sales + certified reseller network",
            "Regions: North America (priority), EMEA (H2 2025)",
            "Pricing: Per-seat SaaS ($29/mo), Enterprise flat-rate ($4,200/yr)",
            "Marketing: Thought leadership, webinars, targeted PPC",
        ],
    )

    # ---- Slide 7: Financial Projections ----
    add_title_content_slide(
        prs,
        title_text="Financial Projections",
        body_lines=[
            "Year 1 ARR target: $3.8M (based on 220 enterprise accounts)",
            "Gross margin: 72% (SaaS model)",
            "Breakeven: Month 14 post-GA launch",
            "Series A target: $12M to fund S&M scale-up",
        ],
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
