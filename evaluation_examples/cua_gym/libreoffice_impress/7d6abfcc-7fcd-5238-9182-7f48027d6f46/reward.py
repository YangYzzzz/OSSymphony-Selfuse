"""
Reward Script: Update slide 5 title text and alignment
Task ID: osworld_impress_title_set_aligned_005
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 5 title text is 'Competitive Landscape' (0.5 points)
  Component 2: Slide 5 title paragraph alignment is CENTER (0.5 points)
"""

import os
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_set_aligned_005'


def persist_app_state():
    """Send Ctrl+S to persist any unsaved GUI edits in LibreOffice Impress."""
    try:
        import pyautogui
        import time
        os.environ["DISPLAY"] = ":0"
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.5)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Task: Rename slide 5 title from 'Competition' to 'Competitive Landscape'
          and set the title paragraph alignment to CENTER.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"CRITICAL: Presentation has {len(prs.slides)} slides, expected at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide5 = prs.slides[4]  # 0-indexed, slide 5 is index 4

    # Find the title shape on slide 5
    title_shape = None
    for shape in slide5.shapes:
        if shape.has_text_frame and shape.name == "Title 1":
            title_shape = shape
            break

    if title_shape is None:
        print("CRITICAL: Title shape ('Title 1') not found on slide 5")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Title text is 'Competitive Landscape' (0.5 points)
    try:
        # Get the text from the first non-empty paragraph in the title
        title_text = None
        title_para = None
        for para in title_shape.text_frame.paragraphs:
            if para.text.strip():
                title_text = para.text.strip()
                title_para = para
                break

        expected_title = "Competitive Landscape"
        if title_text == expected_title:
            print(f"PASS: Component 1 — Title text is '{title_text}' (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 1 — Expected title '{expected_title}', found '{title_text}'")
    except Exception as e:
        print(f"ERROR: Component 1 — Could not check title text: {e}")

    # Component 2: Title paragraph alignment is CENTER (0.5 points)
    try:
        alignment = title_para.alignment if title_para is not None else None

        # PP_ALIGN.CENTER == 2
        if alignment == PP_ALIGN.CENTER:
            print(f"PASS: Component 2 — Title alignment is CENTER ({alignment}) (0.5 pts)")
            total_score += 0.5
        else:
            align_name = str(alignment) if alignment is not None else "LEFT (default/None)"
            print(f"FAIL: Component 2 — Expected CENTER alignment, found {align_name}")
    except Exception as e:
        print(f"ERROR: Component 2 — Could not check title alignment: {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point: test against canonical artifact path on the VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
