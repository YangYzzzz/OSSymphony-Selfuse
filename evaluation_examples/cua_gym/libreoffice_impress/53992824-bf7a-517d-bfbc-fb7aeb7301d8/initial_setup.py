"""
Initial Setup: 15-slide Training presentation with notes on slide 10 only.
Task ID: impress_ndo_017
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_017'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_content(slide, title_text, body_lines):
    """Helper to populate a slide with title and body content."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    # Find the body placeholder (index 1 typically)
    body_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            body_ph = ph
            break
    if body_ph and body_lines:
        tf = body_ph.text_frame
        tf.clear()
        for i, line in enumerate(body_lines):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                p = tf.add_paragraph()
                p.text = line


def create_initial():
    prs = Presentation()

    # Slide layout references
    title_layout = prs.slide_layouts[0]       # Title Slide
    content_layout = prs.slide_layouts[1]     # Title and Content
    blank_layout = prs.slide_layouts[5]       # Blank

    # ── Slide 1: Title Slide ──
    s1 = prs.slides.add_slide(title_layout)
    s1.shapes.title.text = "Employee Training Program 2025"
    s1.placeholders[1].text = "Onboarding & Professional Development\nAcme Corporation"

    # ── Slide 2: Agenda ──
    s2 = prs.slides.add_slide(content_layout)
    add_title_and_content(s2, "Agenda", [
        "1. Welcome & Introductions",
        "2. Company Overview",
        "3. Department Orientation",
        "4. Software Tools & Systems",
        "5. Compliance & Safety",
        "6. Q&A Session",
    ])

    # ── Slide 3: Company Overview ──
    s3 = prs.slides.add_slide(content_layout)
    add_title_and_content(s3, "Company Overview", [
        "Founded in 2003 in San Francisco",
        "Over 2,400 employees across 12 offices globally",
        "Core business: Enterprise logistics solutions",
        "Revenue: $1.2B (FY 2024)",
        "Named Top Employer by Glassdoor three consecutive years",
    ])

    # ── Slide 4: Mission & Values ──
    s4 = prs.slides.add_slide(content_layout)
    add_title_and_content(s4, "Mission & Values", [
        "Mission: Simplify global supply chains",
        "Integrity in every interaction",
        "Innovation through collaboration",
        "Customer success is our success",
        "Sustainability and social responsibility",
    ])

    # ── Slide 5: Department Structure ──
    s5 = prs.slides.add_slide(content_layout)
    add_title_and_content(s5, "Department Structure", [
        "Engineering — Led by VP Sarah Chen",
        "Sales & Marketing — Led by VP Marcus Johnson",
        "Operations — Led by VP Priya Patel",
        "Human Resources — Led by VP David Kim",
        "Finance & Accounting — Led by CFO Rachel Torres",
    ])

    # ── Slide 6: IT Systems Overview ──
    s6 = prs.slides.add_slide(content_layout)
    add_title_and_content(s6, "IT Systems Overview", [
        "Internal portal: AcmeConnect (intranet)",
        "Project management: Jira & Confluence",
        "Communication: Slack and Microsoft Teams",
        "CRM: Salesforce",
        "HR & Payroll: Workday",
    ])

    # ── Slide 7: Security Protocols ──
    s7 = prs.slides.add_slide(content_layout)
    add_title_and_content(s7, "Security Protocols", [
        "Two-factor authentication required for all systems",
        "VPN mandatory when working remotely",
        "Phishing awareness training: quarterly",
        "Report incidents to security@acmecorp.com",
        "Password rotation every 90 days",
    ])

    # ── Slide 8: Compliance Training ──
    s8 = prs.slides.add_slide(content_layout)
    add_title_and_content(s8, "Compliance Training", [
        "Anti-harassment policy review — due within 30 days",
        "Data privacy (GDPR/CCPA) certification",
        "Export control regulations overview",
        "Workplace safety: fire exits, first aid stations",
        "Code of conduct acknowledgment form",
    ])

    # ── Slide 9: Benefits & Resources ──
    s9 = prs.slides.add_slide(content_layout)
    add_title_and_content(s9, "Benefits & Resources", [
        "Health insurance: Aetna PPO / Kaiser HMO",
        "401(k) match up to 6% of salary",
        "Tuition reimbursement: $5,250/year",
        "Employee Assistance Program (EAP)",
        "Gym membership discount: 40% off",
    ])

    # ── Slide 10: Software Workflow Demo ──
    s10 = prs.slides.add_slide(content_layout)
    add_title_and_content(s10, "Software Workflow Demo", [
        "Step 1: Log in to AcmeConnect",
        "Step 2: Navigate to your team dashboard",
        "Step 3: Create a new project request",
        "Step 4: Attach required documentation",
        "Step 5: Submit for manager approval",
    ])
    # Add notes to slide 10 (CRITICAL for the task)
    s10.notes_slide.notes_text_frame.text = "Demonstrate the software workflow step by step."

    # ── Slide 11: Hands-On Exercise ──
    s11 = prs.slides.add_slide(content_layout)
    add_title_and_content(s11, "Hands-On Exercise", [
        "Practice creating a project request in AcmeConnect",
        "Use the test environment (sandbox mode)",
        "Follow the 5-step workflow from the demo",
        "Raise your hand if you need assistance",
        "Target completion time: 15 minutes",
    ])
    # NO notes on slide 11

    # ── Slide 12: Common Troubleshooting ──
    s12 = prs.slides.add_slide(content_layout)
    add_title_and_content(s12, "Common Troubleshooting", [
        "Login issues: Reset via IT Self-Service portal",
        "VPN connection drops: Switch to backup server",
        "Slow dashboard: Clear browser cache and cookies",
        "Missing permissions: Submit access request form",
        "System outage: Check status.acmecorp.com",
    ])
    # NO notes on slide 12

    # ── Slide 13: Advanced Features ──
    s13 = prs.slides.add_slide(content_layout)
    add_title_and_content(s13, "Advanced Features", [
        "Automated report scheduling",
        "Custom dashboard widgets",
        "API integrations with third-party tools",
        "Bulk data import/export via CSV",
        "Role-based access control configuration",
    ])
    # NO notes on slide 13

    # ── Slide 14: Next Steps ──
    s14 = prs.slides.add_slide(content_layout)
    add_title_and_content(s14, "Next Steps", [
        "Complete compliance training modules by April 15",
        "Set up all IT accounts within first week",
        "Schedule 1:1 with your direct manager",
        "Join your department Slack channel",
        "Review the employee handbook on AcmeConnect",
    ])

    # ── Slide 15: Thank You ──
    s15 = prs.slides.add_slide(content_layout)
    add_title_and_content(s15, "Thank You!", [
        "Questions? Contact HR at hr@acmecorp.com",
        "IT Support: helpdesk@acmecorp.com | ext. 4500",
        "Training portal: learn.acmecorp.com",
        "Welcome to the Acme team!",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
