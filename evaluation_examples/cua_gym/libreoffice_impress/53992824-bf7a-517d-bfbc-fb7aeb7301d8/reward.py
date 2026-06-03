"""
Reward Script: Copy notes from slide 10 to slides 11-13, appending ' (continued)'
Task ID: impress_ndo_017
Domain: libreoffice_impress
Scoring:
  Component 1 (0.35): Slide 11 notes == 'Demonstrate the software workflow step by step. (continued)'
  Component 2 (0.35): Slide 12 notes == 'Demonstrate the software workflow step by step. (continued)'
  Component 3 (0.30): Slide 13 notes == 'Demonstrate the software workflow step by step. (continued)'
  Precondition gate: Slide 10 notes must remain unchanged
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_017'
EXPECTED_NOTES = 'Demonstrate the software workflow step by step. (continued)'
SLIDE10_NOTES = 'Demonstrate the software workflow step by step.'


def persist_app_state():
    """Save any unsaved LibreOffice changes before verification."""
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        import time
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def get_notes(slide):
    """Safely extract notes text from a slide."""
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Must have at least 13 slides
    if len(prs.slides) < 13:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 13")
        print("REWARD: 0.0")
        return 0.0

    # Precondition gate: Slide 10 notes must be unchanged
    slide10_notes = get_notes(prs.slides[9])  # 0-indexed
    if slide10_notes != SLIDE10_NOTES:
        print(f"FAIL: Slide 10 notes were modified. Expected: {repr(SLIDE10_NOTES)}, Found: {repr(slide10_notes)}")
        print("REWARD: 0.0")
        return 0.0
    else:
        print(f"GATE: Slide 10 notes unchanged: {repr(slide10_notes)}")

    # Component 1: Slide 11 notes (0.35 points)
    try:
        slide11_notes = get_notes(prs.slides[10])
        if slide11_notes == EXPECTED_NOTES:
            print(f"PASS: Component 1 - Slide 11 notes correct: {repr(slide11_notes)} (0.35 pts)")
            total_score += 0.35
        else:
            print(f"FAIL: Component 1 - Slide 11 notes. Expected: {repr(EXPECTED_NOTES)}, Found: {repr(slide11_notes)}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Slide 12 notes (0.35 points)
    try:
        slide12_notes = get_notes(prs.slides[11])
        if slide12_notes == EXPECTED_NOTES:
            print(f"PASS: Component 2 - Slide 12 notes correct: {repr(slide12_notes)} (0.35 pts)")
            total_score += 0.35
        else:
            print(f"FAIL: Component 2 - Slide 12 notes. Expected: {repr(EXPECTED_NOTES)}, Found: {repr(slide12_notes)}")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Slide 13 notes (0.30 points)
    try:
        slide13_notes = get_notes(prs.slides[12])
        if slide13_notes == EXPECTED_NOTES:
            print(f"PASS: Component 3 - Slide 13 notes correct: {repr(slide13_notes)} (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 3 - Slide 13 notes. Expected: {repr(EXPECTED_NOTES)}, Found: {repr(slide13_notes)}")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
