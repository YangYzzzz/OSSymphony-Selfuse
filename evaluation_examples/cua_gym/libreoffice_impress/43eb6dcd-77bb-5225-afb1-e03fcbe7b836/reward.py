"""
Reward Script: Insert a 6-column by 2-row table on slide 1 with weekly schedule headers
Task ID: impress_tct_024
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Table exists on slide 1 with 6 columns and 2 rows
  Component 2 (0.4): Row 1 contains 'Monday' through 'Saturday' in order
  Component 3 (0.2): Row 2 cells are all empty
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_024'

EXPECTED_DAYS = ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday', 'Saturday']


def find_table_on_slide(slide):
    """Find the first table shape on a slide, return the table or None."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return shape.table
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 1 slide
    if len(prs.slides) < 1:
        print("FAIL: Presentation has no slides")
        print("REWARD: 0.0")
        return 0.0

    slide1 = prs.slides[0]
    table = find_table_on_slide(slide1)

    # Component 1: Table exists on slide 1 with correct dimensions (0.4 points)
    try:
        if table is None:
            print("FAIL: Component 1 -- No table found on slide 1")
        else:
            num_rows = len(table.rows)
            num_cols = len(table.columns)
            if num_rows == 2 and num_cols == 6:
                print(f"PASS: Component 1 -- Table found on slide 1 with {num_rows} rows x {num_cols} cols (0.4 pts)")
                total_score += 0.4
            else:
                print(f"FAIL: Component 1 -- Table dimensions are {num_rows} rows x {num_cols} cols, expected 2 rows x 6 cols")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Row 1 contains day names Monday through Saturday (0.4 points)
    try:
        if table is not None and len(table.rows) >= 1 and len(table.columns) >= 6:
            row1_values = [table.cell(0, c).text.strip() for c in range(6)]
            matches = sum(1 for actual, expected in zip(row1_values, EXPECTED_DAYS)
                         if actual == expected)
            if matches == 6:
                print(f"PASS: Component 2 -- Row 1 contains all correct day names: {row1_values} (0.4 pts)")
                total_score += 0.4
            else:
                print(f"FAIL: Component 2 -- Row 1 values: {row1_values}, expected: {EXPECTED_DAYS} ({matches}/6 match)")
        else:
            print("FAIL: Component 2 -- Table missing or wrong dimensions, cannot check day names")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Row 2 cells are all empty (0.2 points)
    try:
        if table is not None and len(table.rows) >= 2 and len(table.columns) >= 6:
            row2_values = [table.cell(1, c).text.strip() for c in range(6)]
            all_empty = all(v == '' for v in row2_values)
            if all_empty:
                print(f"PASS: Component 3 -- Row 2 is empty as expected (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 3 -- Row 2 has non-empty cells: {row2_values}")
        else:
            print("FAIL: Component 3 -- Table missing or wrong dimensions, cannot check row 2")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
