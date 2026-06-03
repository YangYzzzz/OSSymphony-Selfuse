"""
Initial Setup: Insert a 6-column by 2-row table on slide 1 for weekly schedule headers.
Task ID: impress_tct_024
Domain: libreoffice_impress

Initial state: 3-slide presentation. Slide 1 has title 'This Week' and NO table.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_024'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title slide "This Week" (NO table) ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
    slide1.shapes.title.text = "This Week"
    slide1.placeholders[1].text = "Weekly Planning Overview"

    # --- Slide 2: Team Meeting Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide2.shapes.title.text = "Team Meeting Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Review project milestones for Q2"
    p2 = body2.add_paragraph()
    p2.text = "Discuss resource allocation for new hires"
    p2.level = 0
    p3 = body2.add_paragraph()
    p3.text = "Address client feedback from Sprint 14"
    p3.level = 0
    p4 = body2.add_paragraph()
    p4.text = "Plan team building event for April"
    p4.level = 0

    # --- Slide 3: Key Reminders ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide3.shapes.title.text = "Key Reminders"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Submit timesheets by Friday 5 PM"
    r1 = body3.add_paragraph()
    r1.text = "Office closed on April 18 for maintenance"
    r1.level = 0
    r2 = body3.add_paragraph()
    r2.text = "Quarterly review presentations due April 25"
    r2.level = 0
    r3 = body3.add_paragraph()
    r3.text = "New parking policy takes effect May 1"
    r3.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
