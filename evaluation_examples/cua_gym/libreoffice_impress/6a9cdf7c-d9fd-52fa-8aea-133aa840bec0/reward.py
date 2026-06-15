"""
FINAL REWARD SCRIPT - SUCCESS
Task: I’m wrestling with slide 283 right now—need to drop in a 10-column, 2-row table and tint the very first row with the ‘Gray 10%’ fill (hex #E6E6E6). What’s the fastest way to set that up in LibreOffice Impress?
Generated: 2025-09-10 21:39:00
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation
from pptx.enum.dml import MSO_FILL

# -------------------------------------------------------------
# Reward Script: Verify LibreOffice Impress task completion
# Task: On slide 283 there must be a 10-column × 2-row table
#       whose ENTIRE first row is tinted Gray 10 % (#E6E6E6).
# Scoring (progressive – max 1.0):
#   0.1  Slide 283 exists (non-natural: index very high)
#   0.2  A table exists on that slide
#   0.3  Table dimensions exactly 2 rows × 10 columns
#   0.4  Colour correctness – proportional to the fraction of
#        first-row cells that have the exact solid fill #E6E6E6
# -------------------------------------------------------------

FILE_PATH = (
    "/home/user/"
    "im_wrestling_with_slide_283_right_nowneed_to_drop_in_a_10_column_2_row_table_and_tint_the_very_first_golden.pptx"
)
TARGET_SLIDE_NUMBER = 283  # 1-based index in user request
EXPECTED_ROWS = 2
EXPECTED_COLS = 10
EXPECTED_RGB = (230, 230, 230)  # #E6E6E6 Gray 10 %


# Helper -----------------------------------------------------------------

def _rgb_tuple(color_format):
    """Return (r,g,b) if solid RGB is present, else None."""
    try:
        rgb = color_format.rgb  # returns pptx.dml.color.RGBColor or None
    except Exception:
        return None
    if rgb is None:
        return None
    return (rgb[0], rgb[1], rgb[2])


def _verify_first_row_colour(table):
    """Return (matching_cells, total_first_row_cells)."""
    first_row = table.rows[0]
    matches = 0
    total = len(first_row.cells)
    for cell in first_row.cells:
        fill = cell.fill
        if fill.type != MSO_FILL.SOLID:
            continue
        if _rgb_tuple(fill.fore_color) == EXPECTED_RGB:
            matches += 1
    return matches, total


# Main verification -------------------------------------------------------

def verify_task(file_path: str) -> float:
    print(f"Starting verification for file: {file_path}")
    score = 0.0

    # Check file presence
    if not os.path.exists(file_path):
        print("✗ File not found – task not completed.")
        print("REWARD: 0.0")
        return 0.0

    # Load presentation (PPTX)
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Could not open presentation: {e}")
        print("REWARD: 0.0")
        return 0.0
    print(f"✓ Presentation opened: {len(prs.slides)} slides detected")

    # ------------------------------------------------------------------
    # 1. Slide 283 existence (0.1)
    # ------------------------------------------------------------------
    if len(prs.slides) < TARGET_SLIDE_NUMBER:
        print(f"✗ Slide {TARGET_SLIDE_NUMBER} missing – only {len(prs.slides)} slides.")
        print("REWARD: 0.0")
        return 0.0
    score += 0.1
    slide = prs.slides[TARGET_SLIDE_NUMBER - 1]
    print(f"✓ Slide {TARGET_SLIDE_NUMBER} exists (0.1 points)")

    # ------------------------------------------------------------------
    # 2. Locate table(s) on that slide (0.2)
    # ------------------------------------------------------------------
    table_found = None
    for shape in slide.shapes:
        if shape.has_table:
            table_found = shape.table
            break
    if table_found is None:
        print("✗ No table found on the target slide.")
        print(f"REWARD: {score}")
        return score  # Partial credit possible only for slide existence
    score += 0.2
    print("✓ Table detected on slide (0.2 points)")

    # ------------------------------------------------------------------
    # 3. Verify table dimensions (0.3)
    # ------------------------------------------------------------------
    rows, cols = len(table_found.rows), len(table_found.columns)
    print(f"Table dimensions detected: {rows} rows × {cols} cols")
    if rows == EXPECTED_ROWS and cols == EXPECTED_COLS:
        score += 0.3
        print("✓ Table dimensions match expected 2×10 (0.3 points)")
    else:
        print("✗ Table dimensions incorrect – no dimension points")

    # ------------------------------------------------------------------
    # 4. Verify first-row colour (up to 0.4 proportional)
    # ------------------------------------------------------------------
    matches, total = _verify_first_row_colour(table_found)
    ratio = matches / total if total else 0.0
    colour_points = 0.4 * ratio
    score += colour_points
    if colour_points > 0:
        print(
            f"✓ First-row colour matches in {matches}/{total} cells "
            f"(ratio {ratio:.2f}) – {colour_points:.2f} points"
        )
    else:
        print("✗ First-row cells are not correctly tinted – 0 colour points")

    # ------------------------------------------------------------------
    final_score = round(min(score, 1.0), 2)
    print(f"Total verification score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score


# Execute when run as script ---------------------------------------------
if __name__ == "__main__":
    verify_task(FILE_PATH)

