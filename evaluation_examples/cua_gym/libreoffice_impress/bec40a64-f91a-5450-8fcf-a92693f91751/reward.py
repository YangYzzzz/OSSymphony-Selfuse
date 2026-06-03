"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 235 is sticking out like a sore thumb—the body text is the wrong style. In LibreOffice Impress, how can I switch that slide’s content placeholder to Liberation Serif 18 pt and make sure the paragraph alignment is set to Justified, without altering the rest of the deck?
Generated: 2025-09-10 21:58:31
Status: success
Model: azure-o3
Total Steps: 7
"""

import os
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as PH
from pptx.util import Pt


def verify_slide_235_style(file_path: str) -> float:
    """Verify that on slide 235 the BODY/CONTENT placeholder text
    is set to Liberation Serif 18 pt and paragraphs are Justified.

    Progressive scoring (max 1.0):
      • Font name correct  -> 0.4
      • Font size correct  -> 0.3
      • Alignment justified-> 0.3
    Returns the calculated reward (0.0-1.0) and prints debug details.
    """

    max_score = 1.0
    score = 0.0

    print(f"Verifying task on file: {file_path}")

    # ----- prerequisite checks -----
    if not os.path.exists(file_path):
        print("✗ File not found")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error opening presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    target_idx = 234  # 0-based index for slide 235
    if len(prs.slides) <= target_idx:
        print(f"✗ Presentation has only {len(prs.slides)} slides – slide 235 missing")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[target_idx]

    # Placeholder types that represent the *body/content* area
    body_placeholder_types = {PH.BODY, PH.OBJECT, PH.VERTICAL_BODY, PH.VERTICAL_OBJECT}

    expected_font_name = "Liberation Serif"
    expected_font_size = Pt(18)

    name_ok = True
    size_ok = True
    align_ok = True

    runs_checked = 0
    paragraphs_checked = 0

    # ----- inspect only the content placeholder -----
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        # Only evaluate if it is a placeholder of the body/content family
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            if ph_type not in body_placeholder_types:
                continue  # Skip titles & other placeholders
        else:
            continue  # Skip non-placeholder shapes – they are not the content placeholder

        # Evaluate every non-empty paragraph & its runs
        for paragraph in shape.text_frame.paragraphs:
            if not paragraph.text or not paragraph.text.strip():
                continue  # ignore empty para
            paragraphs_checked += 1

            # Alignment check
            if paragraph.alignment != PP_ALIGN.JUSTIFY:
                align_ok = False

            # Run-level checks
            for run in paragraph.runs:
                if not run.text or not run.text.strip():
                    continue
                runs_checked += 1
                font = run.font

                # Font name check (must be explicitly Liberation Serif)
                if font.name is None or font.name.strip() != expected_font_name:
                    name_ok = False
                # Font size check (must be explicitly 18 pt, allow tiny tolerance)
                if font.size is None or abs(font.size - expected_font_size) > Pt(0.1):
                    size_ok = False

    # ----- guard: ensure we actually inspected something -----
    if paragraphs_checked == 0 or runs_checked == 0:
        print("✗ No content placeholder text found to verify on slide 235")
        print("REWARD: 0.0")
        return 0.0

    # ----- debug summary -----
    print(f"Checked paragraphs: {paragraphs_checked}, runs: {runs_checked}")
    print(f"Font name correct: {'✓' if name_ok else '✗'}")
    print(f"Font size correct: {'✓' if size_ok else '✗'}")
    print(f"Alignment justified: {'✓' if align_ok else '✗'}")

    # ----- progressive scoring -----
    if name_ok:
        score += 0.4
    if size_ok:
        score += 0.3
    if align_ok:
        score += 0.3

    final_score = round(min(score, max_score), 2)
    print(f"REWARD: {final_score}")
    return final_score


# ----------------------------------------------------------------------------------
# Actual call when this script is executed directly (path hard-coded for evaluation)
if __name__ == "__main__":
    FILE_PATH = "/home/user/slide_235_is_sticking_out_like_a_sore_thumbthe_body_text_is_the_wrong_style_in_libreoffice_impress_h_golden.pptx"
    verify_slide_235_style(FILE_PATH)

