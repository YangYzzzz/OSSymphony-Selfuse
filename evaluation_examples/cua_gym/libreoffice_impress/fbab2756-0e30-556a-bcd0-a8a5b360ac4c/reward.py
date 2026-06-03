"""
Reward Script: Duplicate slide 3 and insert the copy after slide 5.
Task ID: osworld_impress_slide_duplication_reorder_003
Domain: libreoffice_impress
Scoring:
  Component 1: Presentation has 8 slides (was 7 before task) — 0.4 pts
  Component 2: Slide 6 is a duplicate of slide 3 (Methodology content) — 0.4 pts
  Component 3: Slides 7-8 are the original slides 6-7 (Future Directions, Conclusion) — 0.2 pts
"""

import os

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_slide_duplication_reorder_003'

# Expected content for Methodology slide (slide 3) — used to verify the duplicate
METHODOLOGY_KEYWORDS = ['Methodology', 'Data Collection', 'Preprocessing Pipeline', 'Modeling Approach', 'Evaluation Metrics']

# Expected titles for shifted slides (original slides 6 and 7 should be at positions 7 and 8 after duplication)
EXPECTED_SLIDE7_TITLE = 'Future Directions'
EXPECTED_SLIDE8_TITLE = 'Conclusion'


def get_all_text(slide):
    """Return all non-empty text strings from a slide."""
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                txt = para.text.strip()
                if txt:
                    texts.append(txt)
    return texts


def get_slide_title(slide):
    """Return the first non-empty text from a slide (used as title)."""
    texts = get_all_text(slide)
    return texts[0] if texts else ''


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide count is exactly 8 (task requires duplicating slide 3, adding one slide) (0.4 points)
    try:
        num_slides = len(prs.slides)
        if num_slides == 8:
            print(f"PASS: Component 1 — Slide count is 8 ({num_slides} slides found) (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 — Expected 8 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — Could not check slide count: {e}")

    # Component 2: Slide 6 (index 5) is a duplicate of slide 3 — contains Methodology content (0.4 points)
    try:
        if len(prs.slides) >= 6:
            slide6 = prs.slides[5]  # 0-indexed, so index 5 = slide 6
            slide6_texts = get_all_text(slide6)
            slide6_text_joined = ' '.join(slide6_texts)

            # Check that all expected Methodology keywords appear in slide 6
            keywords_found = [kw for kw in METHODOLOGY_KEYWORDS if kw in slide6_text_joined]
            if len(keywords_found) == len(METHODOLOGY_KEYWORDS):
                print(f"PASS: Component 2 — Slide 6 contains Methodology content (all {len(METHODOLOGY_KEYWORDS)} keywords found) (0.4 pts)")
                total_score += 0.4
            else:
                missing = [kw for kw in METHODOLOGY_KEYWORDS if kw not in slide6_text_joined]
                print(f"FAIL: Component 2 — Slide 6 missing Methodology keywords: {missing}")
                print(f"  Slide 6 first text: {get_slide_title(slide6)!r}")
        else:
            print(f"FAIL: Component 2 — Not enough slides to check slide 6 (only {len(prs.slides)} slides)")
    except Exception as e:
        print(f"ERROR: Component 2 — Could not check slide 6 content: {e}")

    # Component 3: Slides 7 and 8 are the original slides 6 and 7 (Future Directions, Conclusion) (0.2 points)
    try:
        if len(prs.slides) >= 8:
            slide7_title = get_slide_title(prs.slides[6])  # 0-indexed index 6 = slide 7
            slide8_title = get_slide_title(prs.slides[7])  # 0-indexed index 7 = slide 8

            slide7_ok = EXPECTED_SLIDE7_TITLE.lower() in slide7_title.lower()
            slide8_ok = EXPECTED_SLIDE8_TITLE.lower() in slide8_title.lower()

            if slide7_ok and slide8_ok:
                print(f"PASS: Component 3 — Slides 7 and 8 are correctly shifted: slide7={slide7_title!r}, slide8={slide8_title!r} (0.2 pts)")
                total_score += 0.2
            else:
                if not slide7_ok:
                    print(f"FAIL: Component 3 — Slide 7 expected '{EXPECTED_SLIDE7_TITLE}', found {slide7_title!r}")
                if not slide8_ok:
                    print(f"FAIL: Component 3 — Slide 8 expected '{EXPECTED_SLIDE8_TITLE}', found {slide8_title!r}")
        else:
            print(f"FAIL: Component 3 — Not enough slides to check slides 7-8 (only {len(prs.slides)} slides)")
    except Exception as e:
        print(f"ERROR: Component 3 — Could not check slides 7-8: {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path on the VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
