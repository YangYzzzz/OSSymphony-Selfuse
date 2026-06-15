"""
Initial Setup: Duplicate slide 3 and insert the copy after slide 5.
Task ID: osworld_impress_slide_duplication_reorder_003
Domain: libreoffice_impress

Creates a 7-slide conference presentation. Slide 3 contains Methodology content.
The agent task is to duplicate slide 3 and place the copy after slide 5.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_slide_duplication_reorder_003'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_content_slide(prs, layout_idx, title_text, content_lines):
    """Add a slide with title and content bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    # Set content (placeholder index 1)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            for i, line in enumerate(content_lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.level = 0
            break
    return slide


def create_initial():
    prs = Presentation()

    # Use standard 16:9 widescreen dimensions
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide (Layout 0) ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Advances in Renewable Energy Systems"
    for ph in slide1.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = "International Conference on Sustainable Technology\nDr. Elena Marchetti & Dr. James Okafor\nMarch 2025"
            break

    # --- Slide 2: Overview / Agenda (Layout 1) ---
    add_title_content_slide(prs, 1, "Agenda",
        [
            "1. Introduction & Motivation",
            "2. Literature Review",
            "3. Methodology",
            "4. Experimental Results",
            "5. Discussion",
            "6. Future Directions",
            "7. Conclusion & Questions",
        ]
    )

    # --- Slide 3: Methodology (Layout 1) ---
    add_title_content_slide(prs, 1, "Methodology",
        [
            "Data Collection",
            "  - 5-year dataset from 12 solar farms across 4 climatic zones",
            "  - Hourly irradiance, temperature, and power output readings",
            "Preprocessing Pipeline",
            "  - Outlier removal via IQR filtering (threshold: 1.5x)",
            "  - Missing value imputation using seasonal ARIMA",
            "Modeling Approach",
            "  - Hybrid CNN-LSTM architecture for time-series forecasting",
            "  - Cross-validation: 5-fold stratified by season",
            "Evaluation Metrics",
            "  - RMSE, MAE, R² over held-out test set (2024 Q4)",
        ]
    )

    # --- Slide 4: Experimental Results (Layout 1) ---
    add_title_content_slide(prs, 1, "Experimental Results",
        [
            "Model Performance Summary",
            "  - RMSE: 4.7 kWh (12% improvement over baseline)",
            "  - MAE:  3.2 kWh",
            "  - R²:   0.94",
            "Key Findings",
            "  - CNN-LSTM outperforms SARIMA by 18% on cloudy-day forecasts",
            "  - Feature importance: solar elevation angle (31%), humidity (24%)",
            "Ablation Study",
            "  - Removing temporal attention drops R² by 0.06",
            "  - Data augmentation with synthetic weather improves RMSE by 7%",
        ]
    )

    # --- Slide 5: Discussion (Layout 1) ---
    add_title_content_slide(prs, 1, "Discussion",
        [
            "Interpretation of Results",
            "  - High R² confirms model captures seasonal and diurnal patterns",
            "  - Humid subtropical zones show highest prediction variability",
            "Comparison with Prior Work",
            "  - Surpasses Liu et al. (2023) RMSE by 9% on equivalent dataset",
            "  - Comparable to Zhang & Patel (2024) with 40% less training data",
            "Limitations",
            "  - Limited to ground-mounted utility-scale installations",
            "  - Requires minimum 2 years of historical data for convergence",
        ]
    )

    # --- Slide 6: Future Directions (Layout 1) ---
    add_title_content_slide(prs, 1, "Future Directions",
        [
            "Short-term (6–12 months)",
            "  - Extend to wind turbine output forecasting",
            "  - Integrate satellite imagery as auxiliary input",
            "Medium-term (1–2 years)",
            "  - Deploy real-time inference pipeline at partner utility (GridCo SE)",
            "  - Federated learning across geographically distributed farms",
            "Long-term Vision",
            "  - Unified multi-source renewable forecasting platform",
            "  - Open benchmark dataset release (target: NeurIPS 2026)",
        ]
    )

    # --- Slide 7: Conclusion (Layout 1) ---
    add_title_content_slide(prs, 1, "Conclusion",
        [
            "Summary",
            "  - Proposed CNN-LSTM hybrid achieves state-of-the-art solar forecasting",
            "  - Validated across 12 diverse solar farms over 5 years",
            "Contributions",
            "  - Novel temporal attention mechanism tailored to irradiance cycles",
            "  - Publicly available preprocessing toolkit (GitHub: elena-m/solar-prep)",
            "Acknowledgements",
            "  - Funded by EU Horizon 2020 grant #892341",
            "  - Data provided by SunTrack GmbH and HeliosData Inc.",
            "Questions?",
            "  Contact: e.marchetti@uni-vienna.ac.at",
        ]
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Slide count: {len(prs.slides)} (expected 7)')

    # GUI-ready startup: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
