"""
FINAL REWARD SCRIPT - SUCCESS
Task: Insert a hyperlink 'mailto:support@example.com' on the text 'Contact us'.
Generated: 2025-10-17 16:08:55
Status: success
Model: azure-o3
Total Steps: 1
"""

from pptx import Presentation
import os


def verify_hyperlink(file_path: str) -> float:
    """Verify that the presentation contains the text 'Contact us' and that
    this text (or the shape/run containing it) has a hyperlink pointing to
    'mailto:support@example.com'.  Returns a score between 0.0 and 1.0.
    """
    print(f"Verifying presentation: {file_path}")
    score = 0.0
    max_score = 1.0

    # Task-specific expectations
    target_phrase = "contact us"  # compare case-insensitively
    expected_link = "mailto:support@example.com"

    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0  # nothing to grade if the file is missing

    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded PPTX with {len(prs.slides)} slide(s)")
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        return 0.0

    # Flags we will set while scanning
    phrase_found = False       # did we see the text at all?
    hyperlink_correct = False  # did we see it with the correct hyperlink?

    # Iterate through every slide and every shape
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            # Helper to pull text (if the shape has any)
            shape_text = shape.text.strip() if hasattr(shape, "text") else ""

            # 1) Check for the phrase in the overall shape text
            if target_phrase in shape_text.lower():
                phrase_found = True
                print(f"  ✓ Found text 'Contact us' on slide {slide_idx}, shape {shape_idx}")

            # 2) Check a possible shape-level hyperlink (click_action)
            try:
                if hasattr(shape, "click_action") and shape.click_action is not None:
                    addr = shape.click_action.hyperlink.address
                    if addr and addr.lower() == expected_link:
                        if target_phrase in shape_text.lower():
                            hyperlink_correct = True
                            print(f"  ✓ Correct shape-level hyperlink on slide {slide_idx}, shape {shape_idx}")
            except Exception:
                # Some shapes don't expose click_action – safely ignore
                pass

            # 3) Dive deeper: run-level hyperlinks inside text frames
            if shape.has_text_frame:
                tf = shape.text_frame
                for para_idx, paragraph in enumerate(tf.paragraphs, start=1):
                    # The phrase might be split across runs, so we track contiguous
                    collecting = False
                    collected_text = ""
                    current_addr = None

                    for run_idx, run in enumerate(paragraph.runs, start=1):
                        run_text = run.text or ""
                        # Update phrase_found even if not hyperlinked
                        if target_phrase in run_text.lower():
                            phrase_found = True

                        # Attempt to read any hyperlink on this run
                        run_addr = None
                        try:
                            if run.hyperlink is not None:
                                run_addr = run.hyperlink.address
                        except Exception:
                            pass

                        # If run has the expected hyperlink, collect contiguous text
                        if run_addr and run_addr.lower() == expected_link:
                            if not collecting:
                                collecting = True
                                collected_text = run_text
                                current_addr = run_addr
                            else:
                                collected_text += run_text
                        else:
                            # We left a hyperlinked segment – evaluate what we collected
                            if collecting:
                                if target_phrase in collected_text.lower():
                                    hyperlink_correct = True
                                collecting = False
                                collected_text = ""
                                current_addr = None

                    # Flush at paragraph end if still collecting
                    if collecting and target_phrase in collected_text.lower():
                        hyperlink_correct = True

    # ----- Scoring -----
    if phrase_found:
        print("✓ 'Contact us' text found in presentation (0.5 points)")
        score += 0.5
    else:
        print("✗ 'Contact us' text not found (0 points)")

    if hyperlink_correct:
        print("✓ 'Contact us' text correctly hyperlinked to mailto:support@example.com (0.5 points)")
        score += 0.5
    else:
        print("✗ Correct hyperlink not found on 'Contact us' text (0 points)")

    final_score = min(score, max_score)
    print(f"Total Score: {final_score}/{max_score}")
    return final_score


if __name__ == "__main__":
    PRESENTATION_PATH = "/home/user/insert_a_hyperlink_mailtosupportexamplecom_on_the_text_contact_us.pptx"
    reward_value = verify_hyperlink(PRESENTATION_PATH)
    print(f"REWARD: {reward_value}")
