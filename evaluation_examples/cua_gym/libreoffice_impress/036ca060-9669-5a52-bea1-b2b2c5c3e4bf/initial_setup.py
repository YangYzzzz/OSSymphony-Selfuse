"""
Initial Setup: Create a 5-slide presentation with varied font sizes
Task ID: impstruct_020
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impstruct_020'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_run(paragraph, text, font_name, font_size_pt, bold=False, italic=False, color=None):
    """Add a run with specified formatting."""
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = RGBColor(*color)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ========== SLIDE 1: Title Slide ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    # Title
    txBox = slide1.shapes.add_textbox(Inches(1.5), Inches(1.0), Inches(10), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_text_run(p, "Q4 2025 Design Review", "Arial", 28, bold=True, color=(0x1A, 0x3C, 0x6E))

    # Subtitle
    txBox2 = slide1.shapes.add_textbox(Inches(2.0), Inches(3.0), Inches(9), Inches(1.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    add_text_run(p2, "Product & UX Team - Annual Design Strategy", "Calibri", 20, color=(0x55, 0x55, 0x55))

    # Footnote
    txBox3 = slide1.shapes.add_textbox(Inches(3.0), Inches(6.0), Inches(7), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    add_text_run(p3, "Confidential - Internal Use Only | Last Updated: March 2025", "Calibri", 10, italic=True, color=(0x99, 0x99, 0x99))

    # ========== SLIDE 2: Overview ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])

    # Title
    txBox = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    add_text_run(p, "Project Overview", "Arial", 28, bold=True, color=(0x1A, 0x3C, 0x6E))

    # Subtitle
    txBox2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(10), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    add_text_run(p2, "Key Objectives and Timeline", "Calibri", 20, color=(0x33, 0x66, 0x99))

    # Body text
    txBox3 = slide2.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11), Inches(4.0))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True

    body_items = [
        "Redesign the customer-facing dashboard to improve usability metrics by 35%",
        "Implement responsive layouts across all product pages for mobile-first experience",
        "Consolidate three legacy navigation patterns into a single unified sidebar",
        "Reduce average task completion time from 4.2 minutes to under 2.5 minutes",
        "Launch beta testing program with 200 enterprise customers by end of Q4",
    ]
    for i, item in enumerate(body_items):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        add_text_run(p, f"  {item}", "Georgia", 16, color=(0x33, 0x33, 0x33))

    # Footnote
    txBox4 = slide2.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(8), Inches(0.5))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    add_text_run(p4, "Source: Product Analytics Dashboard, Feb 2025 Report", "Calibri", 10, italic=True, color=(0x99, 0x99, 0x99))

    # ========== SLIDE 3: Research Findings ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])

    txBox = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    add_text_run(p, "User Research Findings", "Arial", 28, bold=True, color=(0x1A, 0x3C, 0x6E))

    txBox2 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(10), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    add_text_run(p2, "Insights from 150+ User Interviews", "Calibri", 20, color=(0x33, 0x66, 0x99))

    txBox3 = slide3.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(5.5), Inches(4.0))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    findings_left = [
        "78% of users found the current navigation confusing when switching between modules",
        "Average onboarding time for new users is 3.8 days, exceeding the 2-day target",
        "Mobile usage increased 62% year-over-year but satisfaction scores dropped 18 points",
    ]
    for i, item in enumerate(findings_left):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        add_text_run(p, f"  {item}", "Georgia", 16, color=(0x33, 0x33, 0x33))

    txBox4 = slide3.shapes.add_textbox(Inches(7.0), Inches(2.5), Inches(5.5), Inches(4.0))
    tf4 = txBox4.text_frame
    tf4.word_wrap = True
    findings_right = [
        "Enterprise customers request role-based dashboard customization most frequently",
        "Integration with Slack and Teams cited as critical by 84% of power users",
        "Data export features used by only 12% but rated essential by those who use them",
    ]
    for i, item in enumerate(findings_right):
        if i == 0:
            p = tf4.paragraphs[0]
        else:
            p = tf4.add_paragraph()
        add_text_run(p, f"  {item}", "Georgia", 16, color=(0x33, 0x33, 0x33))

    txBox5 = slide3.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(8), Inches(0.5))
    tf5 = txBox5.text_frame
    p5 = tf5.paragraphs[0]
    add_text_run(p5, "Methodology: Semi-structured interviews, Oct-Dec 2024", "Calibri", 10, italic=True, color=(0x99, 0x99, 0x99))

    # ========== SLIDE 4: Design Proposals ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])

    txBox = slide4.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    add_text_run(p, "Design Proposals", "Arial", 28, bold=True, color=(0x1A, 0x3C, 0x6E))

    txBox2 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(10), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    add_text_run(p2, "Three Approaches Under Consideration", "Calibri", 20, color=(0x33, 0x66, 0x99))

    proposals = [
        ("Option A: Incremental Refresh", "Minimal disruption to existing workflows. Update color palette, typography, and spacing. Estimated 6-week timeline with 2 developers."),
        ("Option B: Component-Based Redesign", "Replace legacy UI components with a modern design system. Enables consistent experience across features. Estimated 14-week timeline."),
        ("Option C: Full Platform Overhaul", "Ground-up redesign with new information architecture. Highest impact but requires 6-month commitment and dedicated design sprint team."),
    ]
    y_pos = 2.5
    for title, desc in proposals:
        txBox3 = slide4.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(11), Inches(1.2))
        tf3 = txBox3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        add_text_run(p3, title, "Arial", 20, bold=True, color=(0x2A, 0x5C, 0x8E))
        p4 = tf3.add_paragraph()
        add_text_run(p4, desc, "Georgia", 16, color=(0x33, 0x33, 0x33))
        y_pos += 1.3

    txBox6 = slide4.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(8), Inches(0.5))
    tf6 = txBox6.text_frame
    p6 = tf6.paragraphs[0]
    add_text_run(p6, "Cost estimates pending finance review - preliminary figures only", "Calibri", 10, italic=True, color=(0x99, 0x99, 0x99))

    # ========== SLIDE 5: Next Steps & Timeline ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])

    txBox = slide5.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    add_text_run(p, "Next Steps & Timeline", "Arial", 28, bold=True, color=(0x1A, 0x3C, 0x6E))

    txBox2 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(10), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    add_text_run(p2, "Implementation Roadmap", "Calibri", 20, color=(0x33, 0x66, 0x99))

    steps = [
        "Week 1-2: Finalize design direction based on stakeholder feedback",
        "Week 3-4: Create high-fidelity prototypes for core user flows",
        "Week 5-8: Develop and integrate new components into staging environment",
        "Week 9-10: Conduct usability testing with 50 beta participants",
        "Week 11-12: Address feedback, optimize performance, prepare for launch",
    ]
    txBox3 = slide5.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11), Inches(3.5))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    for i, step in enumerate(steps):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        add_text_run(p, f"  {step}", "Georgia", 16, color=(0x33, 0x33, 0x33))

    txBox4 = slide5.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(8), Inches(0.5))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    add_text_run(p4, "Contact: design-team@company.com | Project Slack: #design-review-q4", "Calibri", 10, italic=True, color=(0x99, 0x99, 0x99))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
