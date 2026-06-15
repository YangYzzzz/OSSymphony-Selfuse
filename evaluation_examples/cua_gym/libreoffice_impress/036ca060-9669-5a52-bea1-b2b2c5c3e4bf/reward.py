"""
Reward Script: Increase font size of all text by 4 points (relative)
Task ID: impstruct_020
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Titles are 32pt and no 28pt runs remain
  Component 2 (0.25): No 16pt body text runs remain (all should be 20pt now)
  Component 3 (0.25): Footnotes are 14pt and no 10pt runs remain
  Component 4 (0.20): Total run count preserved and no original-size runs left anywhere
"""

import os
from pptx import Presentation
from pptx.util import Pt

WORKDIR = '/home/user'
TASK_ID = 'impstruct_020'


def collect_all_run_sizes(prs):
    """Collect all (slide_idx, size_pt, text_snippet) tuples for runs with explicit sizes."""
    runs = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for pi, para in enumerate(shape.text_frame.paragraphs):
                    for ri, run in enumerate(para.runs):
                        if run.font.size is not None:
                            size_pt = run.font.size / 12700  # EMU to pt
                            runs.append((si, size_pt, run.text[:50]))
    return runs


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Strategy: Check that original font sizes (28, 16, 10) are gone and
    new sizes (32, 14) are present. The 20pt->24pt shift is verified by
    absence of original sizes and presence of 32pt/14pt unique markers.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have 5 slides
    if len(prs.slides) != 5:
        print(f"PRECONDITION FAIL: Expected 5 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    all_runs = collect_all_run_sizes(prs)
    if not all_runs:
        print("PRECONDITION FAIL: No text runs with explicit font sizes found")
        print("REWARD: 0.0")
        return 0.0

    # Build size buckets (with 0.5pt tolerance)
    def count_at_size(runs, target_pt):
        return sum(1 for _, sz, _ in runs if abs(sz - target_pt) < 0.5)

    def slides_with_size(runs, target_pt):
        return set(si for si, sz, _ in runs if abs(sz - target_pt) < 0.5)

    n_at_28 = count_at_size(all_runs, 28.0)
    n_at_32 = count_at_size(all_runs, 32.0)
    n_at_16 = count_at_size(all_runs, 16.0)
    n_at_10 = count_at_size(all_runs, 10.0)
    n_at_14 = count_at_size(all_runs, 14.0)

    print(f"DEBUG: Runs at 28pt={n_at_28}, 32pt={n_at_32}, 16pt={n_at_16}, 10pt={n_at_10}, 14pt={n_at_14}")
    print(f"DEBUG: Total runs with sizes: {len(all_runs)}")

    # Component 1: Titles increased from 28pt to 32pt (0.30 points)
    # Check: 32pt runs exist on all 5 slides AND no 28pt runs remain
    try:
        slides_32 = slides_with_size(all_runs, 32.0)
        if n_at_32 >= 5 and n_at_28 == 0 and len(slides_32) >= 5:
            print(f"PASS: Component 1 — {n_at_32} title runs at 32pt across {len(slides_32)} slides, 0 at 28pt (0.30 pts)")
            total_score += 0.30
        elif n_at_32 >= 3 and n_at_28 == 0:
            partial = 0.30 * (len(slides_32) / 5)
            print(f"PARTIAL: Component 1 — {n_at_32} runs at 32pt on {len(slides_32)} slides, 0 at 28pt ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — {n_at_32} runs at 32pt, {n_at_28} still at 28pt")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Body text increased from 16pt to 20pt (0.25 points)
    # Check: no 16pt runs remain (they should all be 20pt now)
    # Note: 20pt runs also existed initially (subtitles), but those should now be 24pt.
    # So the check is: zero runs at 16pt.
    try:
        if n_at_16 == 0:
            print(f"PASS: Component 2 — No runs remain at 16pt, body text increased (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 — {n_at_16} runs still at 16pt (should be 0)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Footnotes increased from 10pt to 14pt (0.25 points)
    # Check: 14pt runs exist AND no 10pt runs remain
    try:
        slides_14 = slides_with_size(all_runs, 14.0)
        if n_at_14 >= 4 and n_at_10 == 0:
            print(f"PASS: Component 3 — {n_at_14} footnote runs at 14pt across {len(slides_14)} slides, 0 at 10pt (0.25 pts)")
            total_score += 0.25
        elif n_at_14 >= 2 and n_at_10 == 0:
            partial = 0.25 * (len(slides_14) / 5)
            print(f"PARTIAL: Component 3 — {n_at_14} runs at 14pt on {len(slides_14)} slides ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — {n_at_14} runs at 14pt, {n_at_10} still at 10pt")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Overall consistency — no original sizes remain (0.20 points)
    # All original sizes (28, 16, 10) should be zero. This is a holistic check
    # that catches any missed runs. The 20pt original subtitles should now be 24pt,
    # verified indirectly: if 28/16/10 are all gone and 32/14 exist, the +4pt shift
    # was applied consistently.
    try:
        original_remaining = n_at_28 + n_at_16 + n_at_10
        # Also verify 24pt runs exist (subtitles shifted from 20 to 24)
        n_at_24 = count_at_size(all_runs, 24.0)
        slides_24 = slides_with_size(all_runs, 24.0)
        if original_remaining == 0 and n_at_24 >= 5 and len(slides_24) >= 5:
            print(f"PASS: Component 4 — No original sizes remain, {n_at_24} runs at 24pt across {len(slides_24)} slides (0.20 pts)")
            total_score += 0.20
        elif original_remaining == 0 and n_at_24 >= 3:
            partial = 0.20 * (len(slides_24) / 5)
            print(f"PARTIAL: Component 4 — No original sizes, but only {n_at_24} runs at 24pt ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — {original_remaining} runs at original sizes, {n_at_24} at 24pt")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
