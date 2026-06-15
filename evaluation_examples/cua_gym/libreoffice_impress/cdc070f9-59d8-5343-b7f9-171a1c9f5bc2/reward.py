"""
Reward Script: Convert outline into strategy presentation
Task ID: impress_wf_035
Domain: libreoffice_impress
Scoring:
  C1 (0.10) - 10 slides total
  C2 (0.15) - Slide 2 has 5 pentagon shapes
  C3 (0.20) - Slides 3-7 have correct pillar titles
  C4 (0.15) - Slide 8 has implementation timeline table (Q1-Q4 x pillars)
  C5 (0.10) - Slide 9 has a pie chart
  C6 (0.15) - All slides have Fade transitions
  C7 (0.15) - Color scheme uses #0D47A1 dark blue
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_035'

EXPECTED_PILLARS = [
    'Market Expansion',
    'Digital Transformation',
    'Talent Development',
    'Operational Excellence',
    'Innovation and R&D',
]


def check_transitions(pptx_path, num_slides):
    """Check how many slides have fade transitions."""
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    fade_count = 0
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            for i in range(1, num_slides + 1):
                try:
                    with zf.open(f'ppt/slides/slide{i}.xml') as f:
                        root = ET.parse(f).getroot()
                        tr = root.find('.//p:transition', ns)
                        if tr is not None and tr.find('.//p:fade', ns) is not None:
                            fade_count += 1
                except KeyError:
                    pass
    except Exception as e:
        print(f"ERROR: Transition check failed: {e}")
    return fade_count


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Exactly 10 slides (0.10 points)
    try:
        if num_slides == 10:
            print(f"PASS: Component 1 - Slide count is 10 (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 1 - Expected 10 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Slide 2 has 5 pentagon (homePlate) shapes (0.15 points)
    try:
        if num_slides >= 2:
            slide2 = prs.slides[1]
            pentagon_count = 0
            pentagon_texts = []
            for shape in slide2.shapes:
                # Check if shape is an auto shape with pentagon geometry
                sp_element = shape._element
                prstGeom = sp_element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
                if prstGeom is not None and prstGeom.get('prst') == 'homePlate':
                    pentagon_count += 1
                    if shape.has_text_frame:
                        pentagon_texts.append(shape.text_frame.text.strip())
            if pentagon_count == 5:
                print(f"PASS: Component 2 - Slide 2 has 5 pentagon shapes: {pentagon_texts} (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 2 - Expected 5 pentagons on slide 2, found {pentagon_count}")
        else:
            print(f"FAIL: Component 2 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Slides 3-7 each have correct pillar title (0.20 points)
    # 0.04 points per correct pillar slide
    try:
        pillar_score = 0.0
        if num_slides >= 7:
            for idx, expected_pillar in enumerate(EXPECTED_PILLARS):
                slide = prs.slides[idx + 2]  # slides 3-7 = index 2-6
                # Collect all text from shapes on the slide
                found = False
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if expected_pillar.lower() in text.lower():
                            found = True
                            break
                if found:
                    print(f"PASS: Component 3.{idx+1} - Slide {idx+3} contains '{expected_pillar}' (0.04 pts)")
                    pillar_score += 0.04
                else:
                    print(f"FAIL: Component 3.{idx+1} - Slide {idx+3} missing pillar '{expected_pillar}'")
            total_score += pillar_score
            print(f"  Component 3 subtotal: {pillar_score:.2f}/0.20")
        else:
            print(f"FAIL: Component 3 - Not enough slides for pillars")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Slide 8 has a table with Q1-Q4 columns and pillar rows (0.15 points)
    try:
        if num_slides >= 8:
            slide8 = prs.slides[7]
            table_found = False
            for shape in slide8.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    nrows = len(table.rows)
                    ncols = len(table.columns)

                    # Check header row has Q1-Q4
                    header_texts = [table.cell(0, c).text.strip() for c in range(ncols)]
                    has_quarters = all(q in ' '.join(header_texts) for q in ['Q1', 'Q2', 'Q3', 'Q4'])

                    # Check at least 5 data rows (one per pillar) + header
                    has_pillar_rows = nrows >= 6

                    if has_quarters and has_pillar_rows:
                        print(f"PASS: Component 4 - Slide 8 table: {nrows}x{ncols}, headers={header_texts} (0.15 pts)")
                        total_score += 0.15
                        table_found = True
                    else:
                        print(f"FAIL: Component 4 - Table found but quarters={has_quarters}, rows={nrows} (need >=6)")
                        table_found = True
                    break
            if not table_found:
                print(f"FAIL: Component 4 - No table found on slide 8")
        else:
            print(f"FAIL: Component 4 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Slide 9 has a pie chart (0.10 points)
    try:
        if num_slides >= 9:
            slide9 = prs.slides[8]
            chart_found = False
            for shape in slide9.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    chart = shape.chart
                    # PIE chart type value is 5 in python-pptx
                    chart_type_val = chart.chart_type
                    # Accept any pie variant (PIE=5, PIE_3D=70, PIE_EXPLODED=69, etc.)
                    if chart_type_val in (5, 69, 70, -4102):
                        print(f"PASS: Component 5 - Slide 9 has pie chart (type={chart_type_val}) (0.10 pts)")
                        total_score += 0.10
                        chart_found = True
                    else:
                        print(f"FAIL: Component 5 - Slide 9 chart is type {chart_type_val}, not pie")
                        chart_found = True
                    break
            if not chart_found:
                print(f"FAIL: Component 5 - No chart found on slide 9")
        else:
            print(f"FAIL: Component 5 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    # Component 6: All slides have Fade transitions (0.15 points)
    try:
        if num_slides > 0:
            fade_count = check_transitions(file_path, num_slides)
            if fade_count == num_slides:
                print(f"PASS: Component 6 - All {num_slides} slides have fade transitions (0.15 pts)")
                total_score += 0.15
            elif fade_count > 0:
                partial = 0.15 * fade_count / num_slides
                print(f"PARTIAL: Component 6 - {fade_count}/{num_slides} slides have fade transitions ({partial:.3f} pts)")
                total_score += round(partial, 3)
            else:
                print(f"FAIL: Component 6 - No slides have fade transitions")
        else:
            print(f"FAIL: Component 6 - No slides")
    except Exception as e:
        print(f"ERROR: Component 6 - {e}")

    # Component 7: Color scheme uses #0D47A1 dark blue (0.15 points)
    # Check if at least some slides use 0D47A1 in backgrounds or text
    try:
        dark_blue_found = False
        for slide in prs.slides:
            try:
                bg_fill = slide.background.fill
                if bg_fill.type == 1:  # solid fill
                    rgb_str = str(bg_fill.fore_color.rgb).upper()
                    if rgb_str == '0D47A1':
                        dark_blue_found = True
                        break
            except:
                pass

        if not dark_blue_found:
            # Also check text colors for the dark blue
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                try:
                                    if run.font.color.type is not None:
                                        rgb_str = str(run.font.color.rgb).upper()
                                        if rgb_str == '0D47A1':
                                            dark_blue_found = True
                                            break
                                except:
                                    pass
                                if dark_blue_found:
                                    break
                            if dark_blue_found:
                                break
                    if dark_blue_found:
                        break
                if dark_blue_found:
                    break

        if dark_blue_found:
            print(f"PASS: Component 7 - Dark blue #0D47A1 found in color scheme (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 7 - Dark blue #0D47A1 not found in any background or text")
    except Exception as e:
        print(f"ERROR: Component 7 - {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/Desktop/Strategy_Deck.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
