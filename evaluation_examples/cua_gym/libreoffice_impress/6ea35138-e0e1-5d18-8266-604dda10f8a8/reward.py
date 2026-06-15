"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 217, please insert a 2 × 8 table and set every row to precisely 1.0 cm in height.
Generated: 2025-09-10 19:19:27
Status: success
Model: azure-o3
Total Steps: 3
"""

from pptx import Presentation, util
import os

def verify_task(file_path: str) -> float:
    """Verify that slide 217 contains a 2×8 table whose rows are all 1.0 cm high.

    Returns a progressive score between 0.0-1.0 and prints detailed feedback.
    """

    max_score = 1.0
    score = 0.0

    print(f"Loading presentation: {file_path}")
    if not os.path.exists(file_path):
        print("✗ File not found")
        print("REWARD: 0.0")
        return 0.0

    # ------------------------------------------------------------
    # 1) Load presentation ------------------------------------------------
    # ------------------------------------------------------------
    try:
        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        print(f"✓ Presentation opened (contains {slide_count} slides)")
    except Exception as e:
        print(f"✗ Unable to open PPTX: {e}")
        print("REWARD: 0.0")
        return 0.0

    # ------------------------------------------------------------
    # 2) Check slide 217 exists (index 216) -----------------------
    # ------------------------------------------------------------
    target_idx = 216  # 0-based index
    if slide_count <= target_idx:
        print("✗ Slide 217 does not exist")
        print(f"REWARD: {score}")
        return score  # 0.0 – no further points possible

    print("✓ Slide 217 exists")
    score += 0.4  # presence of the correct slide

    slide = prs.slides[target_idx]

    # ------------------------------------------------------------
    # 3) Find a 2×8 table on that slide ---------------------------
    # ------------------------------------------------------------
    table_found = False
    target_table = None

    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            tbl = shape.table
            if len(tbl.rows) == 8 and len(tbl.columns) == 2:
                table_found = True
                target_table = tbl
                break

    if not table_found:
        print("✗ No 2×8 table found on slide 217")
        print(f"REWARD: {score}")
        return score  # cannot earn further points

    print("✓ Found a 2×8 table on slide 217")
    score += 0.3  # correct table structure

    # ------------------------------------------------------------
    # 4) Verify each row height is exactly 1.0 cm -----------------
    # ------------------------------------------------------------
    expected_height = int(util.Cm(1.0))      # EMU value for 1 cm
    tolerance      = int(util.Cm(0.02))     # ±0.02 cm tolerance (≈2000 EMU)

    all_rows_ok = True
    for idx, row in enumerate(target_table.rows):
        h = row.height
        if h is None:
            print(f"✗ Row {idx+1}: height is not set")
            all_rows_ok = False
            break
        if abs(h - expected_height) > tolerance:
            print(
                f"✗ Row {idx+1}: height {h} EMU differs from 1 cm ({expected_height} EMU) by "
                f"{abs(h - expected_height)} EMU (> tolerance)"
            )
            all_rows_ok = False
            break
    if all_rows_ok:
        print("✓ Every row is precisely 1.0 cm high (within tolerance)")
        score += 0.3
    else:
        print("✗ Not all rows have the required height")

    # ------------------------------------------------------------
    final_score = min(score, max_score)
    print(f"Total score: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score

# ------------------------------------------------------------------
# Execute verification when run as a script -------------------------
# ------------------------------------------------------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/on_slide_217_please_insert_a_2_8_table_and_set_every_row_to_precisely_10_cm_in_height_golden.pptx"
    verify_task(FILE_PATH)
