"""
Initial Setup: Create a 10-slide Annual Report presentation with slide 6 titled 'Revenue Composition' but no chart.
Task ID: impress_exec_021
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_021'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_body_slide(prs, title_text, body_lines, layout_idx=1):
    """Add a slide with title and bulleted body text."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def add_title_only_slide(prs, title_text, layout_idx=5):
    """Add a blank slide with a manual title textbox."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Annual Performance Report 2025"
    slide1.placeholders[1].text = "Prepared by Strategic Finance Division\nConfidential"

    # --- Slide 2: Financial Highlights ---
    add_title_body_slide(prs, "Financial Highlights", [
        "Total Revenue: $58.9M (up 14% YoY)",
        "Gross Margin: 68.2%, improved from 65.8% in 2024",
        "Operating Expenses: $32.1M, controlled growth at 8%",
        "EBITDA: $18.4M with 31.2% margin",
        "Free Cash Flow: $12.7M, strongest quarter in Q4",
        "R&D Investment: $9.3M (15.8% of revenue)",
    ])

    # --- Slide 3: Team & Headcount ---
    add_title_body_slide(prs, "Team & Headcount Overview", [
        "Total Employees: 487 (up from 412 in 2024)",
        "Engineering: 198 (41%), Product: 62 (13%), Sales: 94 (19%)",
        "New Hires: 112 across all departments",
        "Voluntary Attrition: 7.3% (industry avg: 11.2%)",
        "Employee Satisfaction Score: 4.2/5.0",
        "Leadership Development Program: 34 graduates",
    ])

    # --- Slide 4: Operations & Infrastructure ---
    add_title_body_slide(prs, "Operations & Infrastructure", [
        "System Uptime: 99.97% across all production environments",
        "Data Centers: 3 active regions (US-East, EU-West, APAC-Singapore)",
        "Cloud Migration: 82% complete, target 95% by Q2 2026",
        "Security Incidents: 0 critical, 3 moderate (all resolved < 4hrs)",
        "Cost per Transaction: $0.023, down 18% from prior year",
        "Vendor Consolidation: Reduced from 47 to 31 active vendors",
    ])

    # --- Slide 5: Growth Strategy ---
    add_title_body_slide(prs, "Growth Strategy 2026", [
        "Market Expansion: Enter 3 new verticals (Healthcare, Education, Government)",
        "Product Innovation: Launch AI-powered analytics suite in Q2",
        "Partnership Pipeline: 12 strategic partnerships in negotiation",
        "International: Establish EMEA sales office in London by Q3",
        "Customer Success: Target NPS improvement from 62 to 72",
        "M&A: Evaluating 2 acquisition targets in complementary spaces",
    ])

    # --- Slide 6: Revenue Composition (NO CHART - just title and descriptive text) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title
    txBox = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Revenue Composition"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # Descriptive text placeholder
    txBox2 = slide6.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = ("Revenue is categorized into three primary streams: Product sales, "
               "Services engagements, and Licensing agreements. A detailed breakdown "
               "by quarter is needed to visualize the revenue mix trends across 2025.")
    for r in p2.runs:
        r.font.size = Pt(16)
        r.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # --- Slide 7: Marketing Performance ---
    add_title_body_slide(prs, "Marketing Performance", [
        "Marketing Qualified Leads: 4,230 (up 22% YoY)",
        "Cost per Lead: $47.20, down from $58.90",
        "Website Traffic: 1.2M monthly visits, 38% organic",
        "Content Marketing: 48 blog posts, 12 whitepapers, 6 webinars",
        "Social Media Followers: 89K (LinkedIn), 34K (Twitter/X)",
        "Brand Awareness Score: 34% in target market (up from 27%)",
    ])

    # --- Slide 8: Customer Metrics ---
    add_title_body_slide(prs, "Customer Metrics & Retention", [
        "Total Active Customers: 1,847",
        "Enterprise Accounts: 156 (contributing 62% of revenue)",
        "Net Revenue Retention: 118%",
        "Churn Rate: 4.1% annual (down from 5.8%)",
        "Average Contract Value: $31,800 (up 11%)",
        "Customer Support CSAT: 94.2%",
    ])

    # --- Slide 9: 2026 Outlook ---
    add_title_body_slide(prs, "2026 Outlook & Projections", [
        "Revenue Target: $72M (22% growth)",
        "Headcount Plan: Grow to 580 employees",
        "Product Roadmap: 3 major releases planned",
        "Geographic Expansion: EMEA and LATAM focus",
        "Profitability Goal: Achieve 35% EBITDA margin",
        "IPO Readiness: SOX compliance initiative underway",
    ])

    # --- Slide 10: Summary & Next Steps ---
    add_title_body_slide(prs, "Summary & Next Steps", [
        "2025 was a breakout year with strong revenue growth and margin expansion",
        "Key investments in R&D and talent are positioning us for 2026 acceleration",
        "Board approval needed for EMEA office and M&A budget allocation",
        "Next quarterly review: April 15, 2026",
        "Action items distributed via email to all department heads",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
