"""
Reward Script: Add stacked bar chart to slide 6 with revenue data
Task ID: impress_exec_021
Domain: libreoffice_impress
Scoring:
  - Component 1: Chart exists on slide 6 (0.2 pts)
  - Component 2: Chart is stacked bar type (0.15 pts)
  - Component 3: Chart title is 'Revenue Mix by Quarter' (0.15 pts)
  - Component 4: Correct categories Q1-Q4 (0.1 pts)
  - Component 5: Correct series data (Product) (0.15 pts)
  - Component 6: Correct series data (Services) (0.1 pts)
  - Component 7: Correct series data (Licensing) (0.15 pts)
  Total: 1.0
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_021'

# Expected data from the task instruction
EXPECTED_SERIES = {
    'Product': [8.0, 9.5, 11.0, 12.5],
    'Services': [3.0, 3.5, 3.8, 4.2],
    'Licensing': [1.0, 1.5, 1.4, 2.1],
}
EXPECTED_CATEGORIES = ['Q1', 'Q2', 'Q3', 'Q4']
EXPECTED_TITLE = 'Revenue Mix by Quarter'


def values_match(actual, expected, tolerance=0.01):
    """Check if two lists of numbers match within tolerance."""
    if len(actual) != len(expected):
        return False
    return all(abs(a - e) < tolerance for a, e in zip(actual, expected))


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide6 = prs.slides[5]  # 0-indexed

    # Find chart on slide 6
    chart = None
    for shape in slide6.shapes:
        if shape.has_chart:
            chart = shape.chart
            break

    # Component 1: Chart exists on slide 6 (0.2 points)
    try:
        if chart is not None:
            print(f"PASS: Component 1 - Chart found on slide 6 (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 1 - No chart found on slide 6")
            # No chart means no further checks possible
            print(f"\nScore: {total_score}/1.0")
            print(f"REWARD: {total_score}")
            return total_score
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Chart is stacked bar type (0.15 points)
    try:
        from pptx.enum.chart import XL_CHART_TYPE
        # BAR_STACKED (58) or COLUMN_STACKED (52) are both valid stacked bar representations
        chart_type = chart.chart_type
        if chart_type == XL_CHART_TYPE.BAR_STACKED or chart_type == XL_CHART_TYPE.COLUMN_STACKED:
            print(f"PASS: Component 2 - Chart type is stacked bar ({chart_type}) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 - Expected stacked bar chart, found {chart_type}")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Chart title is 'Revenue Mix by Quarter' (0.15 points)
    try:
        if chart.has_title:
            title_text = ''
            if chart.chart_title.has_text_frame:
                title_text = chart.chart_title.text_frame.text.strip()
            if title_text.lower() == EXPECTED_TITLE.lower():
                print(f"PASS: Component 3 - Chart title is '{title_text}' (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 3 - Expected title '{EXPECTED_TITLE}', found '{title_text}'")
        else:
            print(f"FAIL: Component 3 - Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Correct categories Q1-Q4 (0.1 points)
    try:
        actual_cats = [str(c) for c in chart.plots[0].categories]
        if actual_cats == EXPECTED_CATEGORIES:
            print(f"PASS: Component 4 - Categories are {actual_cats} (0.1 pts)")
            total_score += 0.1
        else:
            print(f"FAIL: Component 4 - Expected categories {EXPECTED_CATEGORIES}, found {actual_cats}")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Build series name-to-values map
    series_map = {}
    try:
        from lxml import etree
        ns = {'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'}
        for s in chart.series:
            ser_el = s._element
            tx = ser_el.find('.//c:tx//c:v', ns)
            sname = tx.text.strip() if tx is not None else None
            series_map[sname] = list(s.values)
        print(f"  Series found: {list(series_map.keys())}")
    except Exception as e:
        print(f"ERROR: Could not extract series names: {e}")

    # Component 5: Correct series data - Product (0.15 points)
    try:
        product_vals = series_map.get('Product')
        if product_vals is not None and values_match(product_vals, EXPECTED_SERIES['Product']):
            print(f"PASS: Component 5 - Product series data matches {product_vals} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 5 - Product series expected {EXPECTED_SERIES['Product']}, found {product_vals}")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    # Component 6: Correct series data - Services (0.1 points)
    try:
        services_vals = series_map.get('Services')
        if services_vals is not None and values_match(services_vals, EXPECTED_SERIES['Services']):
            print(f"PASS: Component 6 - Services series data matches {services_vals} (0.1 pts)")
            total_score += 0.1
        else:
            print(f"FAIL: Component 6 - Services series expected {EXPECTED_SERIES['Services']}, found {services_vals}")
    except Exception as e:
        print(f"ERROR: Component 6 - {e}")

    # Component 7: Correct series data - Licensing (0.15 points)
    try:
        licensing_vals = series_map.get('Licensing')
        if licensing_vals is not None and values_match(licensing_vals, EXPECTED_SERIES['Licensing']):
            print(f"PASS: Component 7 - Licensing series data matches {licensing_vals} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 7 - Licensing series expected {EXPECTED_SERIES['Licensing']}, found {licensing_vals}")
    except Exception as e:
        print(f"ERROR: Component 7 - {e}")

    # Round to avoid floating point issues
    total_score = round(total_score, 2)
    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
