"""
Reward Script: Make header row text in table on slide 2 centered horizontally and vertically
Task ID: impress_tct_018
Domain: libreoffice_impress
Scoring:
  Component 1 — Horizontal centering of all 5 header cells (0.5 pts)
  Component 2 — Vertical centering (middle) of all 5 header cells (0.5 pts)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_018'
NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

# Expected number of columns in header row
EXPECTED_COLS = 5


def find_table_on_slide(slide):
    """Find the first table shape on the given slide."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return shape.table
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 2")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: slide 2 must have a table
    slide2 = prs.slides[1]
    table = find_table_on_slide(slide2)
    if table is None:
        print("FAIL: No table found on slide 2")
        print("REWARD: 0.0")
        return 0.0

    num_cols = len(table.columns)
    print(f"INFO: Table has {len(table.rows)} rows x {num_cols} columns")

    # Component 1: Horizontal centering of all header cells (0.5 points)
    # In the initial state, alignment is LEFT. Task requires CENTER.
    try:
        h_centered_count = 0
        for col_idx in range(num_cols):
            cell = table.cell(0, col_idx)
            tf = cell.text_frame
            # Check if every paragraph in this cell is centered
            non_centered = [p for p in tf.paragraphs if p.alignment != PP_ALIGN.CENTER]
            if len(non_centered) == 0:
                h_centered_count += 1
            else:
                print(f"  DETAIL: Cell(0,{col_idx}) para alignment = {non_centered[0].alignment} (expected CENTER)")

        if h_centered_count == num_cols:
            print(f"PASS: Component 1 — All {num_cols} header cells are horizontally centered (0.5 pts)")
            total_score += 0.5
        elif h_centered_count > 0:
            partial = 0.5 * (h_centered_count / num_cols)
            print(f"PARTIAL: Component 1 — {h_centered_count}/{num_cols} header cells horizontally centered ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — No header cells are horizontally centered (0/{num_cols})")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Vertical centering (middle alignment) of all header cells (0.5 points)
    # In the initial state, anchor is "t" (top). Task requires "ctr" (center/middle).
    # Vertical anchor is on the tcPr element, not bodyPr.
    try:
        v_centered_count = 0
        for col_idx in range(num_cols):
            cell = table.cell(0, col_idx)
            tc_el = cell._tc
            # Look for tcPr element and its anchor attribute
            tcPr = tc_el.find(f'{NS}tcPr')
            if tcPr is not None:
                anchor_val = tcPr.get('anchor', 't')  # default is top
            else:
                anchor_val = 't'  # no tcPr means default top

            if anchor_val == 'ctr':
                v_centered_count += 1
            else:
                print(f"  DETAIL: Cell(0,{col_idx}) vertical anchor = '{anchor_val}' (expected 'ctr')")

        if v_centered_count == num_cols:
            print(f"PASS: Component 2 — All {num_cols} header cells are vertically centered (0.5 pts)")
            total_score += 0.5
        elif v_centered_count > 0:
            partial = 0.5 * (v_centered_count / num_cols)
            print(f"PARTIAL: Component 2 — {v_centered_count}/{num_cols} header cells vertically centered ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No header cells are vertically centered (0/{num_cols})")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
