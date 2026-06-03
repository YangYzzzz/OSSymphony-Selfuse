"""
Initial Setup: Employee Directory presentation with left/top-aligned header row in table on slide 2.
Task ID: impress_tct_018
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_018'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Employee Directory"
    slide1.placeholders[1].text = "Acme Corporation — Q1 2025"

    # --- Slide 2: Table slide (5 columns x 10 rows) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Add a title textbox
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Full-Time Employees"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # Table: 10 rows (1 header + 9 data), 5 columns
    rows, cols = 10, 5
    tbl_shape = slide2.shapes.add_table(
        rows, cols,
        Inches(0.5), Inches(1.0),
        Inches(12.0), Inches(5.5)
    )
    table = tbl_shape.table

    headers = ["Employee Name", "Department", "Title", "Start Date", "Office"]
    data = [
        ["Sarah Chen", "Engineering", "Senior Developer", "2023-01-15", "San Francisco"],
        ["Marcus Johnson", "Marketing", "Brand Manager", "2022-06-01", "New York"],
        ["Priya Patel", "Finance", "Financial Analyst", "2024-03-10", "Chicago"],
        ["David Kim", "Engineering", "DevOps Engineer", "2021-11-20", "Seattle"],
        ["Elena Rodriguez", "Human Resources", "HR Coordinator", "2023-08-05", "Austin"],
        ["James O'Brien", "Sales", "Account Executive", "2022-01-30", "Boston"],
        ["Aisha Mohammed", "Product", "Product Manager", "2024-06-15", "San Francisco"],
        ["Lucas Weber", "Engineering", "QA Lead", "2020-09-12", "Denver"],
        ["Mei-Lin Chang", "Operations", "Supply Chain Mgr", "2023-04-22", "Portland"],
    ]

    # Header row — LEFT aligned, TOP aligned (initial state)
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = header
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # Vertical alignment: TOP
        cell.vertical_anchor = MSO_ANCHOR.TOP
        # Header background
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # Data rows
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = Pt(12)
            cell.vertical_anchor = MSO_ANCHOR.TOP
            # Alternate row shading
            if r % 2 == 0:
                cell_fill = cell.fill
                cell_fill.solid()
                cell_fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF4)

    # --- Slide 3: Department Summary ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.6))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Department Breakdown"
    run3 = p3.runs[0]
    run3.font.size = Pt(28)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # Small summary table
    dept_tbl_shape = slide3.shapes.add_table(
        5, 2, Inches(1.0), Inches(1.5), Inches(6.0), Inches(3.0)
    )
    dept_table = dept_tbl_shape.table
    dept_data = [
        ["Department", "Headcount"],
        ["Engineering", "3"],
        ["Marketing", "1"],
        ["Finance", "1"],
        ["Other", "4"],
    ]
    for r, row in enumerate(dept_data):
        for c, val in enumerate(row):
            cell = dept_table.cell(r, c)
            cell.text = val
            if r == 0:
                for run in cell.text_frame.paragraphs[0].runs:
                    run.font.bold = True

    # --- Slide 4: Contact Info ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox4 = slide4.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(5))
    tf4 = txBox4.text_frame
    tf4.word_wrap = True
    p4 = tf4.paragraphs[0]
    p4.text = "Contact Information"
    run4 = p4.runs[0]
    run4.font.size = Pt(28)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    details = [
        "For questions about the employee directory, please reach out to:",
        "",
        "HR Department: hr@acmecorp.com",
        "Office Manager: offices@acmecorp.com",
        "IT Support: helpdesk@acmecorp.com",
        "",
        "Last updated: March 2025",
    ]
    for line in details:
        p_new = tf4.add_paragraph()
        p_new.text = line
        if line:
            for run in p_new.runs:
                run.font.size = Pt(16)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
