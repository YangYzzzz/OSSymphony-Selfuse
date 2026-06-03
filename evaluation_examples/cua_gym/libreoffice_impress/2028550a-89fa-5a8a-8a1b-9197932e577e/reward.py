"""
Reward Script: School Open House Presentation
Task ID: impress_wf_072
Domain: libreoffice_impress
Scoring:
  Component 1: Slide count == 8                              (0.15)
  Component 2: Slide 1 title text contains "Lincoln"         (0.10)
  Component 3: Slide 2 has photo placeholder shape            (0.10)
  Component 4: Slide 3 has schedule table (Time, Activity)    (0.15)
  Component 5: Slide 4 has 6 subject cards                    (0.15)
  Component 6: Slide 6 has calendar table (Date, Event)       (0.10)
  Component 7: Slide 7 has volunteer card shapes              (0.10)
  Component 8: Yellow (#FDD835) header bars present           (0.15)
  Total: 1.0
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_072'


def get_all_text(slide):
    """Get all text from all shapes on a slide."""
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text:
            texts.append(shape.text.strip())
    return texts


def get_shape_fill_colors(slide):
    """Get fill colors (as hex strings) from shapes on a slide."""
    colors = []
    for shape in slide.shapes:
        try:
            fill = shape.fill
            if fill.type is not None:
                try:
                    rgb = str(fill.fore_color.rgb)
                    colors.append(rgb)
                except Exception:
                    pass
        except Exception:
            pass
    return colors


def count_shapes_with_text_matching(slide, targets):
    """Count shapes whose text matches any of the target strings (case-insensitive)."""
    found = set()
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text:
            txt = shape.text.strip().lower()
            for t in targets:
                if t.lower() == txt:
                    found.add(t.lower())
    return found


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Exactly 8 slides (0.15 points)
    try:
        if num_slides == 8:
            print(f"PASS: Component 1 -- Slide count is 8 (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 -- Expected 8 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Slide 1 has "Welcome to Lincoln Elementary" or similar (0.10 points)
    try:
        if num_slides >= 1:
            texts = get_all_text(prs.slides[0])
            all_text = ' '.join(texts).lower()
            if 'lincoln' in all_text and 'welcome' in all_text:
                print(f"PASS: Component 2 -- Slide 1 contains 'Welcome' and 'Lincoln' (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 2 -- Slide 1 text does not contain expected welcome text. Found: {texts[:3]}")
        else:
            print(f"FAIL: Component 2 -- No slides found")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Slide 2 has photo placeholder (circle/rectangle with placeholder text) (0.10 points)
    try:
        if num_slides >= 2:
            slide2 = prs.slides[1]
            texts = get_all_text(slide2)
            all_text = ' '.join(texts).lower()
            has_photo_placeholder = 'photo' in all_text
            has_welcome_or_principal = 'principal' in all_text or 'welcome' in all_text
            if has_photo_placeholder and has_welcome_or_principal:
                print(f"PASS: Component 3 -- Slide 2 has photo placeholder and principal content (0.10 pts)")
                total_score += 0.10
            elif has_photo_placeholder:
                print(f"PARTIAL: Component 3 -- Slide 2 has photo placeholder but missing principal content (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 3 -- Slide 2 missing photo placeholder. Texts: {texts[:3]}")
        else:
            print(f"FAIL: Component 3 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Slide 3 has a schedule table with Time and Activity columns (0.15 points)
    try:
        if num_slides >= 3:
            slide3 = prs.slides[2]
            table_tables = [s for s in slide3.shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE]
            if len(table_tables) > 0:
                table = table_tables[0].table
                headers = [table.cell(0, c).text.strip().lower() for c in range(len(table.columns))] if len(table.rows) >= 2 else []
                if 'time' in headers and 'activity' in headers:
                    print(f"PASS: Component 4 -- Slide 3 has schedule table with Time/Activity headers (0.15 pts)")
                    total_score += 0.15
                elif len(table.rows) >= 2:
                    print(f"PARTIAL: Component 4 -- Slide 3 has a table but missing Time/Activity headers (0.07 pts)")
                    total_score += 0.07
            else:
                print(f"FAIL: Component 4 -- Slide 3 has no table")
        else:
            print(f"FAIL: Component 4 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Slide 4 has 6 subject cards (Math, Reading, Science, Social Studies, Art, PE) (0.15 points)
    try:
        if num_slides >= 4:
            slide4 = prs.slides[3]
            expected_subjects = ['math', 'reading', 'science', 'social studies', 'art', 'pe']
            found_subjects = count_shapes_with_text_matching(slide4, expected_subjects)
            count = len(found_subjects)
            if count >= 6:
                print(f"PASS: Component 5 -- All 6 subject cards found on Slide 4 (0.15 pts)")
                total_score += 0.15
            elif count >= 4:
                pts = round(0.15 * (count / 6), 2)
                print(f"PARTIAL: Component 5 -- {count}/6 subjects found: {found_subjects} ({pts} pts)")
                total_score += pts
            else:
                print(f"FAIL: Component 5 -- Only {count}/6 subjects found: {found_subjects}")
        else:
            print(f"FAIL: Component 5 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # Component 6: Slide 6 has calendar table with Date and Event columns (0.10 points)
    try:
        if num_slides >= 6:
            slide6 = prs.slides[5]
            cal_tables = [s for s in slide6.shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE]
            if len(cal_tables) > 0:
                table = cal_tables[0].table
                headers = [table.cell(0, c).text.strip().lower() for c in range(len(table.columns))] if len(table.rows) >= 2 else []
                if 'date' in headers and 'event' in headers:
                    print(f"PASS: Component 6 -- Slide 6 has calendar table with Date/Event headers (0.10 pts)")
                    total_score += 0.10
                elif len(table.rows) >= 2:
                    print(f"PARTIAL: Component 6 -- Slide 6 has a table but missing Date/Event headers (0.05 pts)")
                    total_score += 0.05
            else:
                print(f"FAIL: Component 6 -- Slide 6 has no table")
        else:
            print(f"FAIL: Component 6 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    # Component 7: Slide 7 has volunteer card shapes (at least 2 volunteer-related shapes) (0.10 points)
    try:
        if num_slides >= 7:
            slide7 = prs.slides[6]
            texts = get_all_text(slide7)
            all_text = ' '.join(texts).lower()
            # Check for volunteer/parent involvement keywords
            volunteer_keywords = ['volunteer', 'helper', 'chaperone', 'coordinator', 'involvement', 'parent']
            keyword_count = sum(1 for kw in volunteer_keywords if kw in all_text)
            # Count rounded rectangle shapes (card shapes) on this slide
            card_shape_count = 0
            for shape in slide7.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and 'rounded' in shape.name.lower():
                    card_shape_count += 1
            if keyword_count >= 2 and card_shape_count >= 2:
                print(f"PASS: Component 7 -- Slide 7 has {card_shape_count} card shapes with volunteer content (0.10 pts)")
                total_score += 0.10
            elif keyword_count >= 2:
                print(f"PARTIAL: Component 7 -- Slide 7 has volunteer content but {card_shape_count} card shapes (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 7 -- Slide 7 missing volunteer content. Keywords found: {keyword_count}")
        else:
            print(f"FAIL: Component 7 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 7 -- {e}")

    # Component 8: Yellow (#FDD835) header bars present on multiple slides (0.15 points)
    try:
        yellow_header_count = 0
        target_color = 'FDD835'
        for i, slide in enumerate(prs.slides):
            colors = get_shape_fill_colors(slide)
            if target_color in colors:
                yellow_header_count += 1
        if yellow_header_count >= 6:
            print(f"PASS: Component 8 -- Yellow (#FDD835) header bars on {yellow_header_count} slides (0.15 pts)")
            total_score += 0.15
        elif yellow_header_count >= 3:
            pts = round(0.15 * (yellow_header_count / 8), 2)
            print(f"PARTIAL: Component 8 -- Yellow headers on {yellow_header_count}/8 slides ({pts} pts)")
            total_score += pts
        else:
            print(f"FAIL: Component 8 -- Yellow (#FDD835) headers found on only {yellow_header_count} slides")
    except Exception as e:
        print(f"ERROR: Component 8 -- {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice
def persist_app_state(domain):
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
