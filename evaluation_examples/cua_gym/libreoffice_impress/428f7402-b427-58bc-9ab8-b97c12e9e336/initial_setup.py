"""
Initial Setup: Conference talk presentation with 10 slides.
Slide 2 has presenter notes in default font (not Arial, not 14pt).
Task ID: impress_tm_086
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_086'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=None):
    """Helper to add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)
    return txBox


def set_notes(slide, notes_text):
    """Set notes for a slide using default formatting (no explicit font)."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = notes_text


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Building Scalable Microservices"
    slide1.placeholders[1].text = "Sarah Chen - Senior Architect\nCloudNative Summit 2025"
    set_notes(slide1, "Welcome everyone. Introduce yourself and mention your 8 years of experience with distributed systems.")

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1]
    tf2 = body2.text_frame
    tf2.clear()
    items = [
        "1. Why Microservices?",
        "2. Architecture Patterns",
        "3. Service Discovery & Load Balancing",
        "4. Data Management Strategies",
        "5. Monitoring & Observability",
        "6. Case Study: Migration at TechCorp",
        "7. Q&A"
    ]
    for i, item in enumerate(items):
        if i == 0:
            tf2.paragraphs[0].text = item
        else:
            p = tf2.add_paragraph()
            p.text = item

    # Slide 2 notes - default font, NOT Arial, NOT 14pt
    set_notes(slide2, "Walk through the agenda quickly. Emphasize that the case study at the end covers a real-world migration from a monolith serving 2 million daily active users. Mention that architecture patterns section will be the longest at roughly 15 minutes.")

    # --- Slide 3: Why Microservices ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Why Microservices?"
    body3 = slide3.placeholders[1]
    tf3 = body3.text_frame
    tf3.clear()
    points3 = [
        "Independent deployment cycles",
        "Technology diversity per service",
        "Fault isolation and resilience",
        "Scalability at the component level",
        "Team autonomy and ownership"
    ]
    for i, pt in enumerate(points3):
        if i == 0:
            tf3.paragraphs[0].text = pt
        else:
            p = tf3.add_paragraph()
            p.text = pt
    set_notes(slide3, "Discuss the monolith pain points first. Reference the 2024 CNCF survey showing 78% of enterprises adopting microservices.")

    # --- Slide 4: Architecture Patterns ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Core Architecture Patterns"
    body4 = slide4.placeholders[1]
    tf4 = body4.text_frame
    tf4.clear()
    points4 = [
        "API Gateway Pattern",
        "Service Mesh (Istio, Linkerd)",
        "Event-Driven Architecture",
        "CQRS and Event Sourcing",
        "Saga Pattern for Distributed Transactions"
    ]
    for i, pt in enumerate(points4):
        if i == 0:
            tf4.paragraphs[0].text = pt
        else:
            p = tf4.add_paragraph()
            p.text = pt
    set_notes(slide4, "Spend extra time on the Saga pattern - most attendees find this the most challenging concept to implement correctly.")

    # --- Slide 5: Service Discovery ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Service Discovery & Load Balancing"
    body5 = slide5.placeholders[1]
    tf5 = body5.text_frame
    tf5.clear()
    points5 = [
        "Client-side vs Server-side Discovery",
        "Consul, Eureka, and etcd comparison",
        "Health checks and circuit breakers",
        "Round-robin vs Weighted load balancing"
    ]
    for i, pt in enumerate(points5):
        if i == 0:
            tf5.paragraphs[0].text = pt
        else:
            p = tf5.add_paragraph()
            p.text = pt
    set_notes(slide5, "Demo the Consul dashboard if time permits. Show the health check failure scenario.")

    # --- Slide 6: Data Management ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Data Management Strategies"
    body6 = slide6.placeholders[1]
    tf6 = body6.text_frame
    tf6.clear()
    points6 = [
        "Database per Service pattern",
        "Shared database anti-pattern",
        "Eventual consistency trade-offs",
        "CDC with Debezium for data sync"
    ]
    for i, pt in enumerate(points6):
        if i == 0:
            tf6.paragraphs[0].text = pt
        else:
            p = tf6.add_paragraph()
            p.text = pt
    set_notes(slide6, "Highlight that shared databases are the number one cause of coupling in failed microservice migrations.")

    # --- Slide 7: Monitoring ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Monitoring & Observability"
    body7 = slide7.placeholders[1]
    tf7 = body7.text_frame
    tf7.clear()
    points7 = [
        "Three Pillars: Logs, Metrics, Traces",
        "OpenTelemetry for instrumentation",
        "Prometheus + Grafana dashboards",
        "Distributed tracing with Jaeger",
        "Alerting best practices"
    ]
    for i, pt in enumerate(points7):
        if i == 0:
            tf7.paragraphs[0].text = pt
        else:
            p = tf7.add_paragraph()
            p.text = pt
    set_notes(slide7, "Show the Grafana dashboard screenshot. Mention the 99.9th percentile latency target of 200ms.")

    # --- Slide 8: Case Study ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Case Study: TechCorp Migration"
    body8 = slide8.placeholders[1]
    tf8 = body8.text_frame
    tf8.clear()
    points8 = [
        "Legacy: 1.2M lines Java monolith (8 years old)",
        "Timeline: 18 months, 4 phases",
        "Result: 23 services, 99.95% uptime",
        "Deployment frequency: weekly to 50x per day",
        "Infrastructure cost reduced by 34%"
    ]
    for i, pt in enumerate(points8):
        if i == 0:
            tf8.paragraphs[0].text = pt
        else:
            p = tf8.add_paragraph()
            p.text = pt
    set_notes(slide8, "This is the core of the talk. Walk through each migration phase in detail. Mention the rollback strategy used in Phase 2.")

    # --- Slide 9: Lessons Learned ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Key Lessons Learned"
    body9 = slide9.placeholders[1]
    tf9 = body9.text_frame
    tf9.clear()
    points9 = [
        "Start with the Strangler Fig pattern",
        "Invest in CI/CD early - it pays off 10x",
        "Contract testing prevents integration failures",
        "Team topology matters as much as architecture",
        "Don't microservice everything - right-size your services"
    ]
    for i, pt in enumerate(points9):
        if i == 0:
            tf9.paragraphs[0].text = pt
        else:
            p = tf9.add_paragraph()
            p.text = pt
    set_notes(slide9, "Emphasize the last point - over-decomposition is the second most common mistake after shared databases.")

    # --- Slide 10: Q&A ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[0])
    slide10.shapes.title.text = "Questions & Discussion"
    slide10.placeholders[1].text = "sarah.chen@techcorp.io\n@sarahchen_arch\ngithub.com/sarahchen/microservices-patterns"
    set_notes(slide10, "Open the floor for questions. If no questions, prompt with the common monolith-to-microservices question.")

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
