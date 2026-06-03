"""
Reward Script: Format presenter notes on slide 2 to 14pt Arial font
Task ID: impress_tm_086
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): Font size is exactly 14pt on all runs of slide 2 notes
  Component 2 (0.3): Font name is explicitly Arial on all runs of slide 2 notes
  Component 3 (0.2): Notes text content is preserved (not empty/corrupted)
  Note: Components 2 and 3 are gated on Component 1 to avoid scoring pre-existing state
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_086'

EXPECTED_NOTES_SUBSTRING = "Walk through the agenda quickly"


def persist_app_state(domain: str):
    """Save any unsaved GUI edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(1.0)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
    except ImportError:
        print("CRITICAL: python-pptx not available")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"CRITICAL: Presentation has {len(prs.slides)} slides, need at least 2")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[1]  # slide 2 (0-indexed)

    # Precondition: slide 2 must have notes with content
    try:
        ns = slide.notes_slide
        tf = ns.notes_text_frame
        notes_text = tf.text.strip()
    except Exception as e:
        print(f"CRITICAL: Cannot access slide 2 notes: {e}")
        print("REWARD: 0.0")
        return 0.0

    if not notes_text:
        print("CRITICAL: Slide 2 has no notes text")
        print("REWARD: 0.0")
        return 0.0

    # Collect all non-empty runs from notes paragraphs
    all_runs = []
    for para in tf.paragraphs:
        for run in para.runs:
            if (run.text or "").strip():
                all_runs.append(run)

    if not all_runs:
        print("CRITICAL: No non-empty runs found in slide 2 notes")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Slide 2 notes has {len(all_runs)} non-empty run(s)")

    # Component 1: Font size is exactly 14pt (177800 EMU) on all runs (0.5 points)
    # This is the PRIMARY discriminator - initial has default size (None or ~20pt), golden has 14pt
    size_correct_count = 0
    try:
        expected_size = 177800  # 14pt in EMU
        for run in all_runs:
            fsize = run.font.size
            if fsize is not None and fsize == expected_size:
                size_correct_count += 1
            else:
                actual_pt = fsize / 12700 if fsize else None
                print(f"FAIL detail: Run '{run.text[:40]}...' has size={fsize} ({actual_pt}pt), expected {expected_size} (14pt)")

        if size_correct_count == len(all_runs):
            print(f"PASS: Component 1 -- All {size_correct_count} run(s) have 14pt font size (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 1 -- {size_correct_count}/{len(all_runs)} runs have 14pt font size")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Font name is Arial AND size is 14pt (0.3 points)
    # Gated on size check to avoid scoring pre-existing Arial font name
    if size_correct_count == len(all_runs):
        try:
            arial_count = 0
            for run in all_runs:
                fname = run.font.name
                if fname and fname.lower() == "arial":
                    arial_count += 1
                else:
                    print(f"FAIL detail: Run '{run.text[:40]}...' has font name={fname}, expected Arial")

            if arial_count == len(all_runs):
                print(f"PASS: Component 2 -- All {arial_count} run(s) have Arial font (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 2 -- {arial_count}/{len(all_runs)} runs have Arial font name")
        except Exception as e:
            print(f"ERROR: Component 2 -- {e}")
    else:
        print("SKIP: Component 2 -- gated on Component 1 (font size must be 14pt first)")

    # Component 3: Notes text content preserved AND correctly formatted (0.2 points)
    # Gated on size check to avoid scoring pre-existing content
    if size_correct_count == len(all_runs):
        try:
            if EXPECTED_NOTES_SUBSTRING.lower() in notes_text.lower():
                print(f"PASS: Component 3 -- Notes content preserved with expected text (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 3 -- Notes text does not contain expected substring")
                print(f"  Expected substring: '{EXPECTED_NOTES_SUBSTRING}'")
                print(f"  Actual text: '{notes_text[:100]}...'")
        except Exception as e:
            print(f"ERROR: Component 3 -- {e}")
    else:
        print("SKIP: Component 3 -- gated on Component 1 (font size must be 14pt first)")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
