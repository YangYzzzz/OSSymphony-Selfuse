"""
Initial Setup: Add semi-transparent dark overlay with bold white text on slide 2
Task ID: impress_sales_046
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFilter
import random

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_046'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_busy_office_image(path, width=1280, height=720):
    """Create a synthetic busy office background image."""
    img = Image.new('RGB', (width, height), (200, 195, 185))
    draw = ImageDraw.Draw(img)

    # Floor
    draw.rectangle([0, height // 2, width, height], fill=(160, 140, 120))

    # Ceiling
    draw.rectangle([0, 0, width, height // 6], fill=(230, 225, 220))

    # Draw desks and furniture shapes
    random.seed(42)
    # Desks
    desk_colors = [(140, 110, 80), (130, 100, 70), (150, 120, 90)]
    for x in range(50, width - 150, 200):
        for y_base in [height // 2 - 60, height // 2 + 80]:
            c = random.choice(desk_colors)
            draw.rectangle([x, y_base, x + 160, y_base + 50], fill=c)
            # Monitor on desk
            draw.rectangle([x + 50, y_base - 60, x + 120, y_base], fill=(30, 30, 35))
            draw.rectangle([x + 55, y_base - 55, x + 115, y_base - 5], fill=(70, 130, 180))

    # Chairs (circles)
    for x in range(100, width - 100, 200):
        for y_base in [height // 2 - 10, height // 2 + 140]:
            c = (60, 60, 65)
            draw.ellipse([x, y_base, x + 40, y_base + 40], fill=c)

    # People silhouettes (simple rectangles + circles for heads)
    people_positions = [(120, 280), (320, 290), (520, 275), (720, 285),
                        (920, 280), (200, 420), (400, 430), (600, 425),
                        (800, 420), (1000, 430)]
    for px, py in people_positions:
        shirt_color = random.choice([(40, 80, 140), (140, 40, 40), (40, 120, 60),
                                      (100, 60, 120), (60, 60, 60)])
        # Body
        draw.rectangle([px - 15, py, px + 15, py + 50], fill=shirt_color)
        # Head
        skin = random.choice([(220, 190, 160), (180, 140, 110), (140, 100, 70)])
        draw.ellipse([px - 12, py - 25, px + 12, py], fill=skin)

    # Windows on back wall
    for wx in range(80, width - 80, 250):
        draw.rectangle([wx, 60, wx + 180, height // 2 - 80], fill=(180, 210, 240))
        # Window frame
        draw.line([wx + 90, 60, wx + 90, height // 2 - 80], fill=(200, 200, 200), width=2)
        draw.line([wx, (60 + height // 2 - 80) // 2, wx + 180, (60 + height // 2 - 80) // 2],
                  fill=(200, 200, 200), width=2)

    # Slight blur for realism
    img = img.filter(ImageFilter.GaussianBlur(radius=1.5))

    img.save(path)
    return path


def create_initial():
    prs = Presentation()
    # Standard 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Impact Pitch 2025"
    if slide1.placeholders[1]:
        slide1.placeholders[1].text = "Transforming Enterprise Efficiency"
    bg1 = slide1.background.fill
    bg1.solid()
    bg1.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)
    # Make title text white
    for run in slide1.shapes.title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.size = Pt(44)
        run.font.bold = True
    for run in slide1.placeholders[1].text_frame.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        run.font.size = Pt(24)

    # --- Slide 2: Background image of busy office (NO overlay, NO problem text) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank or Title Only
    # Create and add background image
    img_path = f'{WORKDIR}/office_bg.png'
    create_busy_office_image(img_path, width=1280, height=720)
    # Add as full-bleed image
    slide2.shapes.add_picture(img_path, 0, 0, slide_w, slide_h)
    # Remove any title placeholder text if present
    for shape in slide2.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""

    # --- Slide 3: Market Opportunity ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    bg3 = slide3.background.fill
    bg3.solid()
    bg3.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)
    txb = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1))
    p = txb.text_frame.paragraphs[0]
    p.text = "Market Opportunity"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Bullet points
    txb2 = slide3.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(10), Inches(4))
    tf = txb2.text_frame
    tf.word_wrap = True
    bullets = [
        "Enterprise productivity software market growing at 12% CAGR",
        "85% of Fortune 500 companies report workflow inefficiencies",
        "Average employee loses 2.5 hours daily to manual processes",
        "Remote work acceleration has doubled digital tool adoption",
    ]
    for i, text in enumerate(bullets):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.text = text
        para.space_after = Pt(14)
        for r in para.runs:
            r.font.size = Pt(20)
            r.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    # --- Slide 4: Our Solution ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    bg4 = slide4.background.fill
    bg4.solid()
    bg4.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)
    txb = slide4.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1))
    p = txb.text_frame.paragraphs[0]
    p.text = "Our Solution"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    txb2 = slide4.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(10), Inches(4))
    tf = txb2.text_frame
    tf.word_wrap = True
    points = [
        "AI-powered workflow automation platform",
        "Seamless integration with 200+ enterprise tools",
        "No-code interface for process design",
        "Real-time analytics and optimization engine",
    ]
    for i, text in enumerate(points):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.text = text
        para.space_after = Pt(14)
        for r in para.runs:
            r.font.size = Pt(20)
            r.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    # --- Slide 5: Traction & Metrics ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    bg5 = slide5.background.fill
    bg5.solid()
    bg5.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)
    txb = slide5.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1))
    p = txb.text_frame.paragraphs[0]
    p.text = "Traction & Metrics"
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Metric boxes
    metrics = [("$4.2M", "ARR"), ("340+", "Enterprise Clients"), ("98%", "Retention Rate")]
    for i, (val, label) in enumerate(metrics):
        x = Inches(1.0 + i * 3.8)
        txm = slide5.shapes.add_textbox(x, Inches(2.5), Inches(3), Inches(2.5))
        tf = txm.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = val
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(48)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0x00, 0xB4, 0xD8)
        p2 = tf.add_paragraph()
        p2.text = label
        p2.alignment = PP_ALIGN.CENTER
        for r in p2.runs:
            r.font.size = Pt(18)
            r.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)

    # --- Slide 6: Business Model ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    bg6 = slide6.background.fill
    bg6.solid()
    bg6.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)
    txb = slide6.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1))
    p = txb.text_frame.paragraphs[0]
    p.text = "Business Model"
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    txb2 = slide6.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(10), Inches(4))
    tf = txb2.text_frame
    tf.word_wrap = True
    lines = [
        "SaaS subscription: $49/user/month (Team), $99/user/month (Enterprise)",
        "Implementation services: one-time setup fees $5K-$50K",
        "Average contract value: $180,000/year",
        "Gross margin: 82%",
        "LTV/CAC ratio: 5.2x",
    ]
    for i, text in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.text = text
        para.space_after = Pt(12)
        for r in para.runs:
            r.font.size = Pt(18)
            r.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    # --- Slide 7: Team ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    bg7 = slide7.background.fill
    bg7.solid()
    bg7.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)
    txb = slide7.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1))
    p = txb.text_frame.paragraphs[0]
    p.text = "Leadership Team"
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    team = [
        ("Elena Martinez", "CEO", "Ex-VP Product at Salesforce, Stanford MBA"),
        ("David Park", "CTO", "Former Principal Engineer at Google, MIT CS PhD"),
        ("Sarah Chen", "COO", "McKinsey alum, scaled ops at Stripe from 50 to 500"),
    ]
    for i, (name, title, bio) in enumerate(team):
        y = Inches(2.0 + i * 1.5)
        txb = slide7.shapes.add_textbox(Inches(0.8), y, Inches(10), Inches(1.2))
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{name} — {title}"
        for r in p.runs:
            r.font.size = Pt(22)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p2 = tf.add_paragraph()
        p2.text = bio
        for r in p2.runs:
            r.font.size = Pt(16)
            r.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)

    # --- Slide 8: The Ask ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    bg8 = slide8.background.fill
    bg8.solid()
    bg8.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)
    txb = slide8.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1))
    p = txb.text_frame.paragraphs[0]
    p.text = "The Ask"
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    txb2 = slide8.shapes.add_textbox(Inches(2), Inches(2.5), Inches(8), Inches(3))
    tf = txb2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Series A: $15M"
    p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.size = Pt(40)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0x00, 0xB4, 0xD8)
    items = [
        "Engineering team expansion (40%)",
        "Go-to-market acceleration (35%)",
        "Infrastructure & security (25%)",
    ]
    for text in items:
        p2 = tf.add_paragraph()
        p2.text = text
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(10)
        for r in p2.runs:
            r.font.size = Pt(20)
            r.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
