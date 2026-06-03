"""
Reward Script: Add semi-transparent dark overlay rectangle with bold white text on slide 2
Task ID: impress_sales_046
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): Black overlay rectangle exists on slide 2
  Component 2 (0.3): Rectangle has ~40% opacity and covers full slide
  Component 3 (0.4): White bold 44pt centered text 'The $50 Billion Problem' on slide 2
"""

import os
import re
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_046'


def persist_app_state(domain):
    """Save any unsaved LibreOffice edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for", domain)
    except Exception as e:
        print("PERSIST_WARN: save hook failed:", e)


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print("CRITICAL: Cannot load file {}: {}".format(file_path, e))
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 2 slides
    if len(prs.slides) < 2:
        print("FAIL: Presentation has fewer than 2 slides")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Find overlay rectangles and text shapes on slide 2
    overlay_rect = None
    text_shape_with_target = None

    for shape in slide2.shapes:
        # Look for auto-shape rectangles (not pictures, not text boxes without fill)
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # Check if it has a solid fill that is black
            try:
                fill = shape.fill
                if fill.type is not None:
                    shape_xml = shape._element
                    xml_str = ET.tostring(shape_xml, encoding='unicode')
                    # Check for black color (000000)
                    if re.search(r'srgbClr\s+val="000000"', xml_str) or re.search(r'srgbClr\s+val="000000"', xml_str.replace("'", '"')):
                        overlay_rect = shape
            except Exception:
                pass

        # Look for text containing "50 Billion Problem"
        if shape.has_text_frame:
            full_text = ""
            for para in shape.text_frame.paragraphs:
                full_text += para.text
            if "50 Billion Problem" in full_text:
                text_shape_with_target = shape

    # Component 1: Black overlay rectangle exists on slide 2 (0.3 points)
    try:
        if overlay_rect is not None:
            print("PASS: Component 1 -- Black overlay rectangle found on slide 2 (name: {}) (0.3 pts)".format(overlay_rect.name))
            total_score += 0.3
        else:
            print("FAIL: Component 1 -- No black overlay rectangle found on slide 2")
    except Exception as e:
        print("ERROR: Component 1 -- {}".format(e))

    # Component 2: Rectangle has ~40% opacity and covers full slide (0.3 points)
    try:
        if overlay_rect is not None:
            comp2_score = 0.0

            # Check opacity (~40%, which is alpha val ~40000 in OOXML, meaning 40% opaque)
            shape_xml = overlay_rect._element
            xml_str = ET.tostring(shape_xml, encoding='unicode')
            alpha_matches = re.findall(r'alpha\s+val="(\d+)"', xml_str)

            if alpha_matches:
                alpha_val = int(alpha_matches[0])
                # alpha val in OOXML is in 1/1000 of percent. 40000 = 40%
                # Allow tolerance: 30000-50000 (30%-50% opacity)
                if 30000 <= alpha_val <= 50000:
                    print("PASS: Component 2a -- Rectangle opacity is {}% (target ~40%)".format(alpha_val / 1000))
                    total_score += 0.15
                else:
                    print("FAIL: Component 2a -- Rectangle opacity is {}%, expected ~40%".format(alpha_val / 1000))
            else:
                print("FAIL: Component 2a -- No alpha/opacity found on rectangle")

            # Check rectangle covers full slide (within 5% tolerance)
            rect_w = overlay_rect.width
            rect_h = overlay_rect.height
            w_ratio = rect_w / slide_width if slide_width > 0 else 0
            h_ratio = rect_h / slide_height if slide_height > 0 else 0

            if w_ratio >= 0.95 and h_ratio >= 0.95:
                print("PASS: Component 2b -- Rectangle covers full slide (w={:.1%}, h={:.1%})".format(w_ratio, h_ratio))
                total_score += 0.15
            else:
                print("FAIL: Component 2b -- Rectangle does not cover full slide (w={:.1%}, h={:.1%})".format(w_ratio, h_ratio))
        else:
            print("FAIL: Component 2 -- No overlay rectangle to check dimensions/opacity")
    except Exception as e:
        print("ERROR: Component 2 -- {}".format(e))

    # Component 3: White bold 44pt centered text 'The $50 Billion Problem' (0.4 points)
    try:
        if text_shape_with_target is not None:
            # Check text content
            full_text = ""
            for para in text_shape_with_target.text_frame.paragraphs:
                full_text += para.text

            if "50 Billion Problem" in full_text:
                print("PASS: Component 3a -- Text contains '50 Billion Problem': {}".format(repr(full_text.strip())))
                total_score += 0.1
            else:
                print("FAIL: Component 3a -- Text '50 Billion Problem' not found")

            # Check font properties on runs containing the target text
            target_runs = []
            target_paras = []
            for para in text_shape_with_target.text_frame.paragraphs:
                for run in para.runs:
                    if run.text and "50 Billion" in run.text:
                        target_runs.append(run)
                        if para not in target_paras:
                            target_paras.append(para)

            if target_runs:
                run = target_runs[0]

                # Check bold
                if run.font.bold:
                    print("PASS: Component 3b -- Text is bold")
                    total_score += 0.05
                else:
                    print("FAIL: Component 3b -- Text is not bold (got {})".format(run.font.bold))

                # Check font size (~44pt = 558800 EMU)
                font_size = run.font.size
                if font_size is not None:
                    pt_size = font_size / 12700
                    if 40 <= pt_size <= 48:
                        print("PASS: Component 3c -- Font size is {:.0f}pt (target 44pt)".format(pt_size))
                        total_score += 0.05
                    else:
                        print("FAIL: Component 3c -- Font size is {:.0f}pt, expected ~44pt".format(pt_size))
                else:
                    print("FAIL: Component 3c -- Font size is None")

                # Check white color
                try:
                    color_rgb = str(run.font.color.rgb)
                    if color_rgb.upper() == "FFFFFF":
                        print("PASS: Component 3d -- Text color is white (FFFFFF)")
                        total_score += 0.1
                    else:
                        print("FAIL: Component 3d -- Text color is {}, expected FFFFFF".format(color_rgb))
                except Exception:
                    print("FAIL: Component 3d -- Cannot read text color")

                # Check centered alignment
                if target_paras:
                    from pptx.enum.text import PP_ALIGN
                    if target_paras[0].alignment == PP_ALIGN.CENTER:
                        print("PASS: Component 3e -- Text is centered")
                        total_score += 0.1
                    else:
                        print("FAIL: Component 3e -- Text alignment is {}, expected CENTER".format(target_paras[0].alignment))
            else:
                print("FAIL: Component 3b-e -- No runs found containing target text")
        else:
            print("FAIL: Component 3 -- No text shape with '50 Billion Problem' found on slide 2")
    except Exception as e:
        print("ERROR: Component 3 -- {}".format(e))

    final_score = min(total_score, 1.0)
    print()
    print("Score: {}/1.0".format(total_score))
    print("REWARD: {}".format(final_score))
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = '{}/{}.pptx'.format(WORKDIR, TASK_ID)
if not os.path.exists(file_path):
    print("File not found: {}".format(file_path))
    print("REWARD: 0.0")
else:
    verify_task(file_path)
