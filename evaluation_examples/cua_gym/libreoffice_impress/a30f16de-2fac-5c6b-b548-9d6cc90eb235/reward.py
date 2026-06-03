"""
Reward Script: Apply pale yellow background to slides 1 and 5 (speaker photo slides)
Task ID: osworld_impress_conditional_bg_image_011
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 1 has pale yellow (FFFFCC) solid background — 0.40 pts
  Component 2: Slide 5 has pale yellow (FFFFCC) solid background — 0.40 pts
  Component 3 (compound): Slide 1 and Slide 5 both changed to FFFFCC AND slides 2,3,4,6 remain FFFFFF — 0.20 pts
    (This guards against accidentally modifying the wrong slides)
"""

import os
from pptx import Presentation
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_conditional_bg_image_011'

# Pale yellow target color
PALE_YELLOW = 'FFFFCC'
WHITE = 'FFFFFF'


def get_slide_bg_rgb(slide):
    """Return the slide background fill RGB hex string, or None if not a solid fill."""
    fill = slide.background.fill
    if fill.type == 1:  # SOLID
        try:
            return str(fill.fore_color.rgb)
        except Exception:
            return None
    return None


def verify_task(file_path):
    """
    Verify that slides 1 and 5 have pale yellow (FFFFCC) backgrounds,
    and that slides 2, 3, 4, 6 remain white (FFFFFF).
    Returns a float between 0.0 and 1.0.
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Verify the presentation has 6 slides
    num_slides = len(prs.slides)
    if num_slides != 6:
        print(f"CRITICAL: Expected 6 slides, found {num_slides}. Presentation structure is wrong.")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide 1 has pale yellow (FFFFCC) background (0.40 points)
    # This FAILS on initial (FFFFFF) → PASSES on golden (FFFFCC) ✓
    slide1_ok = False
    try:
        slide1_bg = get_slide_bg_rgb(prs.slides[0])
        if slide1_bg and slide1_bg.upper() == PALE_YELLOW:
            print(f"PASS: Component 1 — Slide 1 background is pale yellow (FFFFCC) (0.40 pts)")
            total_score += 0.40
            slide1_ok = True
        else:
            print(f"FAIL: Component 1 — Slide 1 background expected FFFFCC, found: {slide1_bg}")
    except Exception as e:
        print(f"ERROR: Component 1 — Could not check slide 1 background: {e}")

    # Component 2: Slide 5 has pale yellow (FFFFCC) background (0.40 points)
    # This FAILS on initial (FFFFFF) → PASSES on golden (FFFFCC) ✓
    slide5_ok = False
    try:
        slide5_bg = get_slide_bg_rgb(prs.slides[4])
        if slide5_bg and slide5_bg.upper() == PALE_YELLOW:
            print(f"PASS: Component 2 — Slide 5 background is pale yellow (FFFFCC) (0.40 pts)")
            total_score += 0.40
            slide5_ok = True
        else:
            print(f"FAIL: Component 2 — Slide 5 background expected FFFFCC, found: {slide5_bg}")
    except Exception as e:
        print(f"ERROR: Component 2 — Could not check slide 5 background: {e}")

    # Component 3 (compound): Slides 1 and 5 changed to FFFFCC AND slides 2,3,4,6 remain white (0.20 points)
    # This is a compound check: it requires BOTH the task changes (slides 1+5 are yellow) AND
    # no over-application (slides 2,3,4,6 are still white).
    # The compound nature ensures this FAILS on initial because slide1_ok and slide5_ok are False.
    try:
        non_target_indices = [1, 2, 3, 5]  # 0-based indices for slides 2, 3, 4, 6
        non_target_white = True
        non_target_failures = []
        for idx in non_target_indices:
            slide_bg = get_slide_bg_rgb(prs.slides[idx])
            if slide_bg is None or slide_bg.upper() != WHITE:
                non_target_white = False
                non_target_failures.append(f"Slide {idx + 1}: {slide_bg}")

        # Compound check: requires both target slides changed AND non-target slides unchanged
        if slide1_ok and slide5_ok and non_target_white:
            print(f"PASS: Component 3 — Selective application correct: only slides 1 and 5 changed to FFFFCC, "
                  f"slides 2,3,4,6 remain FFFFFF (0.20 pts)")
            total_score += 0.20
        elif not slide1_ok or not slide5_ok:
            print(f"FAIL: Component 3 — Target slides not fully changed to FFFFCC yet")
        else:
            print(f"FAIL: Component 3 — Non-target slides have unexpected backgrounds: {non_target_failures}")
    except Exception as e:
        print(f"ERROR: Component 3 — Could not check slide selectivity: {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
