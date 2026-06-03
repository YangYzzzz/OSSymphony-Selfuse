"""
Initial Setup: 6-slide conference keynote deck — all white backgrounds.
Task ID: osworld_impress_conditional_bg_image_011
Domain: libreoffice_impress

Creates a realistic 6-slide conference keynote presentation.
Slides 1 and 5 contain speaker photos (as placeholder rectangles with captions).
ALL slide backgrounds are white — no pale yellow applied.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_conditional_bg_image_011'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_slide_bg_white(slide):
    """Set the slide background to solid white."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                color=None, alignment=PP_ALIGN.LEFT):
    """Helper: add a styled text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_speaker_photo_placeholder(slide, caption_text):
    """Add a grey rectangle to simulate a speaker photo + caption below it."""
    # Grey rectangle to represent speaker photo
    from pptx.util import Inches
    photo_left = Inches(3.0)
    photo_top = Inches(1.5)
    photo_w = Inches(3.5)
    photo_h = Inches(3.5)

    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        photo_left, photo_top, photo_w, photo_h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    shape.line.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    shape.line.width = Pt(1)

    # Photo label inside rectangle
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[Speaker Photo]"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    run.font.italic = True

    # Caption below photo
    caption_top = photo_top + photo_h + Inches(0.1)
    add_textbox(
        slide,
        photo_left, caption_top, photo_w, Inches(0.6),
        caption_text,
        font_size=12,
        bold=False,
        color=RGBColor(0x44, 0x44, 0x44),
        alignment=PP_ALIGN.CENTER,
    )


def create_initial():
    prs = Presentation()
    # Standard widescreen 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    DARK_BLUE = RGBColor(0x1F, 0x3B, 0x6B)
    MEDIUM_GRAY = RGBColor(0x55, 0x55, 0x55)
    LIGHT_GRAY = RGBColor(0xDD, 0xDD, 0xDD)
    BLACK = RGBColor(0x00, 0x00, 0x00)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    blank_layout = prs.slide_layouts[6]  # blank layout

    # ------------------------------------------------------------------
    # Slide 1 — Title / Opening with Speaker Photo
    # ------------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg_white(slide1)

    # Conference title header
    add_textbox(
        slide1,
        Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.8),
        "TechForward Global Summit 2025",
        font_size=13, bold=False,
        color=MEDIUM_GRAY,
        alignment=PP_ALIGN.CENTER,
    )

    # Main title
    add_textbox(
        slide1,
        Inches(0.4), Inches(1.0), Inches(9.2), Inches(1.2),
        "The Future of AI-Driven Enterprise Automation",
        font_size=28, bold=True,
        color=DARK_BLUE,
        alignment=PP_ALIGN.CENTER,
    )

    # Speaker photo placeholder (slide 1 features speaker photo)
    create_speaker_photo_placeholder(slide1, "Dr. Elena Vasquez\nChief AI Officer, NovaTech Inc.")

    # Event details
    add_textbox(
        slide1,
        Inches(0.4), Inches(6.4), Inches(9.2), Inches(0.7),
        "March 18, 2025 | San Francisco Convention Center | Hall B",
        font_size=11, bold=False,
        color=MEDIUM_GRAY,
        alignment=PP_ALIGN.CENTER,
    )

    # ------------------------------------------------------------------
    # Slide 2 — Agenda
    # ------------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_bg_white(slide2)

    add_textbox(
        slide2,
        Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.9),
        "Today's Agenda",
        font_size=30, bold=True,
        color=DARK_BLUE,
        alignment=PP_ALIGN.LEFT,
    )

    # Horizontal rule simulation (thin rectangle)
    hr = slide2.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(9.0), Emu(36000))
    hr.fill.solid()
    hr.fill.fore_color.rgb = DARK_BLUE
    hr.line.fill.background()

    agenda_items = [
        "1.  Opening Remarks & Speaker Introduction",
        "2.  The State of Enterprise AI — 2025 Landscape Report",
        "3.  Case Study: Autonomous Supply Chain at Scale",
        "4.  Panel Discussion: Governance, Ethics & Risk",
        "5.  Live Demo: NovaTech AutoFlow Platform",
        "6.  Q&A and Closing Thoughts",
    ]
    for i, item in enumerate(agenda_items):
        add_textbox(
            slide2,
            Inches(0.8), Inches(1.4 + i * 0.85), Inches(8.5), Inches(0.75),
            item,
            font_size=16, bold=False,
            color=BLACK,
            alignment=PP_ALIGN.LEFT,
        )

    # ------------------------------------------------------------------
    # Slide 3 — Key Statistics
    # ------------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_bg_white(slide3)

    add_textbox(
        slide3,
        Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.9),
        "Enterprise AI Adoption: Key Metrics 2024–2025",
        font_size=26, bold=True,
        color=DARK_BLUE,
        alignment=PP_ALIGN.LEFT,
    )

    hr3 = slide3.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(9.0), Emu(36000))
    hr3.fill.solid()
    hr3.fill.fore_color.rgb = DARK_BLUE
    hr3.line.fill.background()

    stats = [
        ("78%", "of Fortune 500 firms deployed generative AI tools in production"),
        ("$4.2B", "average enterprise AI investment in 2024 (up 3.1× from 2022)"),
        ("12.4M", "knowledge workers augmented by AI copilots globally"),
        ("43%", "reduction in manual data-processing time reported by early adopters"),
        ("2.8×", "productivity multiplier observed in pilot programs"),
    ]
    for i, (stat_val, stat_desc) in enumerate(stats):
        add_textbox(
            slide3,
            Inches(0.8), Inches(1.5 + i * 0.95), Inches(1.5), Inches(0.75),
            stat_val,
            font_size=22, bold=True,
            color=RGBColor(0x0D, 0x6E, 0xBF),
            alignment=PP_ALIGN.RIGHT,
        )
        add_textbox(
            slide3,
            Inches(2.5), Inches(1.55 + i * 0.95), Inches(7.0), Inches(0.65),
            stat_desc,
            font_size=15, bold=False,
            color=BLACK,
            alignment=PP_ALIGN.LEFT,
        )

    # ------------------------------------------------------------------
    # Slide 4 — Case Study
    # ------------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_bg_white(slide4)

    add_textbox(
        slide4,
        Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.9),
        "Case Study: Autonomous Supply Chain at GlobalShip Logistics",
        font_size=22, bold=True,
        color=DARK_BLUE,
        alignment=PP_ALIGN.LEFT,
    )

    hr4 = slide4.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(9.0), Emu(36000))
    hr4.fill.solid()
    hr4.fill.fore_color.rgb = DARK_BLUE
    hr4.line.fill.background()

    case_text = (
        "Challenge:  GlobalShip processed 850,000 shipments monthly with 14% "
        "exception rate, requiring 340 FTE analysts to manage exceptions manually.\n\n"
        "Solution:  Deployed NovaTech AutoFlow v3 across 12 regional hubs, integrating "
        "real-time IoT telemetry, predictive disruption models, and autonomous rerouting.\n\n"
        "Results (18-month post-deployment):\n"
        "   •  Exception rate reduced to 3.1% (78% improvement)\n"
        "   •  Analyst headcount reallocated: 280 FTEs moved to strategic planning roles\n"
        "   •  On-time delivery improved from 82% to 97.3%\n"
        "   •  Annual cost savings: $47.2M across the network"
    )
    add_textbox(
        slide4,
        Inches(0.6), Inches(1.4), Inches(8.8), Inches(5.5),
        case_text,
        font_size=14, bold=False,
        color=BLACK,
        alignment=PP_ALIGN.LEFT,
    )

    # ------------------------------------------------------------------
    # Slide 5 — Speaker Bio / Closing with Speaker Photo
    # ------------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_bg_white(slide5)

    add_textbox(
        slide5,
        Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.85),
        "Meet Your Speaker",
        font_size=28, bold=True,
        color=DARK_BLUE,
        alignment=PP_ALIGN.LEFT,
    )

    hr5 = slide5.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(9.0), Emu(36000))
    hr5.fill.solid()
    hr5.fill.fore_color.rgb = DARK_BLUE
    hr5.line.fill.background()

    # Speaker photo on slide 5
    create_speaker_photo_placeholder(slide5, "Dr. Elena Vasquez")

    bio_text = (
        "Dr. Elena Vasquez is the Chief AI Officer at NovaTech Inc., leading a "
        "450-person team across research, applied ML, and enterprise deployment.\n\n"
        "She holds a PhD in Distributed Systems from MIT and an MBA from INSEAD. "
        "Dr. Vasquez previously led the AI Infrastructure division at Google Cloud "
        "and was founding CTO of Meridian Labs (acquired 2021).\n\n"
        "She is a frequent speaker at NeurIPS, ICML, and Davos, and serves on "
        "the AI Safety Board of the Partnership on AI."
    )
    add_textbox(
        slide5,
        Inches(0.5), Inches(1.4), Inches(2.4), Inches(5.5),
        bio_text,
        font_size=12, bold=False,
        color=BLACK,
        alignment=PP_ALIGN.LEFT,
    )

    # ------------------------------------------------------------------
    # Slide 6 — Thank You / Q&A
    # ------------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_bg_white(slide6)

    add_textbox(
        slide6,
        Inches(0.4), Inches(2.0), Inches(9.2), Inches(1.4),
        "Thank You",
        font_size=52, bold=True,
        color=DARK_BLUE,
        alignment=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide6,
        Inches(0.4), Inches(3.5), Inches(9.2), Inches(0.8),
        "Questions & Discussion",
        font_size=24, bold=False,
        color=MEDIUM_GRAY,
        alignment=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide6,
        Inches(0.4), Inches(5.0), Inches(9.2), Inches(0.6),
        "elena.vasquez@novatech.ai  |  linkedin.com/in/elenavasquez  |  @elenavasquez_ai",
        font_size=13, bold=False,
        color=RGBColor(0x0D, 0x6E, 0xBF),
        alignment=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide6,
        Inches(0.4), Inches(6.3), Inches(9.2), Inches(0.6),
        "Slides available at: novatech.ai/summit2025",
        font_size=11, bold=False,
        color=MEDIUM_GRAY,
        alignment=PP_ALIGN.CENTER,
    )

    # ------------------------------------------------------------------
    # Save
    # ------------------------------------------------------------------
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
