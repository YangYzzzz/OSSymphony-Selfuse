"""
Initial Setup: Physics Demo presentation with Ball shape on slide 1
Task ID: impress_ma_064
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_064'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Only with Ball shape ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide1.shapes.title.text = 'Physics Demonstration'
    # Style the title
    for run in slide1.shapes.title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    # Add subtitle textbox
    txBox = slide1.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = 'Projectile Motion Analysis'
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(18)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # Add Ball (circle shape) at left center of slide
    ball_size = Inches(1.2)
    ball_left = Inches(0.8)
    ball_top = Emu(int((prs.slide_height - ball_size) / 2))
    ball = slide1.shapes.add_shape(MSO_SHAPE.OVAL, ball_left, ball_top, ball_size, ball_size)
    ball.name = 'Ball'
    fill = ball.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xE8, 0x4D, 0x1A)  # Orange-red ball
    # Add outline
    ball.line.color.rgb = RGBColor(0xA0, 0x30, 0x10)
    ball.line.width = Pt(2)

    # Add a ground line shape
    ground = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(6.2), Inches(10), Inches(0.05)
    )
    ground.name = 'Ground'
    ground.fill.solid()
    ground.fill.fore_color.rgb = RGBColor(0x4A, 0x7C, 0x2E)
    ground.line.fill.background()

    # --- Slide 2: Projectile Motion Theory ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide2.shapes.title.text = 'Projectile Motion Equations'
    body2 = slide2.placeholders[1].text_frame
    body2.text = 'Horizontal displacement: x(t) = v₀ · cos(θ) · t'
    p2 = body2.add_paragraph()
    p2.text = 'Vertical displacement: y(t) = v₀ · sin(θ) · t - ½g·t²'
    p2.level = 0
    p3 = body2.add_paragraph()
    p3.text = 'Where v₀ is initial velocity, θ is launch angle'
    p3.level = 1
    p4 = body2.add_paragraph()
    p4.text = 'g = 9.81 m/s² (gravitational acceleration)'
    p4.level = 1
    p5 = body2.add_paragraph()
    p5.text = 'Time of flight: T = 2·v₀·sin(θ) / g'
    p5.level = 0
    p6 = body2.add_paragraph()
    p6.text = 'Maximum height: H = v₀²·sin²(θ) / (2g)'
    p6.level = 0

    # --- Slide 3: Energy Conservation ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide3.shapes.title.text = 'Energy Conservation in Motion'
    body3 = slide3.placeholders[1].text_frame
    body3.text = 'Kinetic Energy: KE = ½mv²'
    pa = body3.add_paragraph()
    pa.text = 'Potential Energy: PE = mgh'
    pa.level = 0
    pb = body3.add_paragraph()
    pb.text = 'Total Mechanical Energy: E = KE + PE = constant'
    pb.level = 0
    pc = body3.add_paragraph()
    pc.text = 'At highest point: KE is minimum, PE is maximum'
    pc.level = 1
    pd = body3.add_paragraph()
    pd.text = 'At lowest point: KE is maximum, PE is minimum'
    pd.level = 1
    pe = body3.add_paragraph()
    pe.text = 'Energy transforms between kinetic and potential forms'
    pe.level = 0

    # --- Slide 4: Experimental Parameters ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide4.shapes.title.text = 'Experimental Setup & Results'
    body4 = slide4.placeholders[1].text_frame
    body4.text = 'Ball mass: 0.25 kg'
    rows = [
        ('Initial velocity: 12.5 m/s', 0),
        ('Launch angle: 45°', 0),
        ('Measured range: 15.9 m', 0),
        ('Theoretical range: 15.94 m', 0),
        ('Percentage error: 0.25%', 1),
        ('Air resistance was negligible in this setup', 1),
    ]
    for text, level in rows:
        p = body4.add_paragraph()
        p.text = text
        p.level = level

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
