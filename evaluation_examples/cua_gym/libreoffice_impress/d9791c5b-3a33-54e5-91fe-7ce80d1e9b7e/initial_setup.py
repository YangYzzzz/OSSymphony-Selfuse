"""
Initial Setup: Create a 5-slide Impress presentation with shapes and text on white backgrounds.
Task ID: impress_el_042
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_el_042'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard 10x7.5 slide dimensions
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    # ====================== SLIDE 1: Title & Logo Shapes ======================
    slide1 = prs.slides.add_slide(blank_layout)
    # White background (explicit)
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Title text box
    txBox = slide1.shapes.add_textbox(Inches(1.5), Inches(1.0), Inches(7), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Overlay Graphics Dashboard"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    # Decorative rectangle
    rect = slide1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(3.0), Inches(6.0), Inches(2.5)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x3A, 0x7C, 0xBD)
    rect.line.fill.background()
    tf2 = rect.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Q1 2025 Strategic Overview"
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.runs[0]
    r2.font.size = Pt(22)
    r2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    r2.font.bold = True

    # Small accent circle
    circle = slide1.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.5), Inches(6.0), Inches(1.0), Inches(1.0)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0xE8, 0x4D, 0x3D)
    circle.line.fill.background()

    # ====================== SLIDE 2: Team Members ======================
    slide2 = prs.slides.add_slide(blank_layout)
    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Section header
    header2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.0), Inches(1.0))
    tf_h2 = header2.text_frame
    p_h2 = tf_h2.paragraphs[0]
    p_h2.text = "Core Team Members"
    p_h2.alignment = PP_ALIGN.LEFT
    r_h2 = p_h2.runs[0]
    r_h2.font.name = "Arial"
    r_h2.font.size = Pt(28)
    r_h2.font.bold = True
    r_h2.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # Team member cards (colored rectangles with text)
    members = [
        ("Sarah Chen", "Lead Designer", RGBColor(0x27, 0xAE, 0x60)),
        ("Marcus Johnson", "Product Manager", RGBColor(0x29, 0x80, 0xB9)),
        ("Elena Rodriguez", "Data Analyst", RGBColor(0x8E, 0x44, 0xAD)),
    ]
    for i, (name, role, color) in enumerate(members):
        card = slide2.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8 + i * 3.0), Inches(2.0), Inches(2.6), Inches(3.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()
        tf_card = card.text_frame
        tf_card.word_wrap = True
        p_name = tf_card.paragraphs[0]
        p_name.text = name
        p_name.alignment = PP_ALIGN.CENTER
        r_name = p_name.runs[0]
        r_name.font.size = Pt(18)
        r_name.font.bold = True
        r_name.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p_role = tf_card.add_paragraph()
        p_role.text = role
        p_role.alignment = PP_ALIGN.CENTER
        r_role = p_role.runs[0]
        r_role.font.size = Pt(14)
        r_role.font.color.rgb = RGBColor(0xEC, 0xF0, 0xF1)

    # ====================== SLIDE 3: KPI Metrics ======================
    slide3 = prs.slides.add_slide(blank_layout)
    fill3 = slide3.background.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    header3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.0), Inches(1.0))
    tf_h3 = header3.text_frame
    p_h3 = tf_h3.paragraphs[0]
    p_h3.text = "Key Performance Indicators"
    r_h3 = p_h3.runs[0]
    r_h3.font.name = "Arial"
    r_h3.font.size = Pt(28)
    r_h3.font.bold = True
    r_h3.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # KPI boxes
    kpis = [
        ("Revenue", "$4.2M", RGBColor(0x16, 0xA0, 0x85)),
        ("Users", "128K", RGBColor(0xF3, 0x9C, 0x12)),
        ("NPS Score", "72", RGBColor(0xE7, 0x4C, 0x3C)),
        ("Retention", "89%", RGBColor(0x34, 0x98, 0xDB)),
    ]
    for i, (label, value, color) in enumerate(kpis):
        box = slide3.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5 + i * 2.3), Inches(2.2), Inches(2.0), Inches(2.0)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf_kpi = box.text_frame
        tf_kpi.word_wrap = True
        p_val = tf_kpi.paragraphs[0]
        p_val.text = value
        p_val.alignment = PP_ALIGN.CENTER
        r_val = p_val.runs[0]
        r_val.font.size = Pt(32)
        r_val.font.bold = True
        r_val.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p_lbl = tf_kpi.add_paragraph()
        p_lbl.text = label
        p_lbl.alignment = PP_ALIGN.CENTER
        r_lbl = p_lbl.runs[0]
        r_lbl.font.size = Pt(14)
        r_lbl.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Accent arrow shape
    arrow = slide3.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(3.5), Inches(5.0), Inches(3.0), Inches(1.0)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0x2E, 0xCC, 0x71)
    arrow.line.fill.background()
    tf_arrow = arrow.text_frame
    p_arrow = tf_arrow.paragraphs[0]
    p_arrow.text = "Trending Up"
    p_arrow.alignment = PP_ALIGN.CENTER
    r_arrow = p_arrow.runs[0]
    r_arrow.font.size = Pt(16)
    r_arrow.font.bold = True
    r_arrow.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # ====================== SLIDE 4: Timeline ======================
    slide4 = prs.slides.add_slide(blank_layout)
    fill4 = slide4.background.fill
    fill4.solid()
    fill4.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    header4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(8.0), Inches(0.8))
    tf_h4 = header4.text_frame
    p_h4 = tf_h4.paragraphs[0]
    p_h4.text = "Project Timeline"
    r_h4 = p_h4.runs[0]
    r_h4.font.name = "Arial"
    r_h4.font.size = Pt(28)
    r_h4.font.bold = True
    r_h4.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # Horizontal timeline bar
    bar = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.5), Inches(8.0), Inches(0.3)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
    bar.line.fill.background()

    # Milestone markers
    milestones = [
        ("Jan 2025", "Research Phase", Inches(1.2)),
        ("Mar 2025", "Prototype", Inches(3.4)),
        ("Jun 2025", "Beta Launch", Inches(5.6)),
        ("Sep 2025", "General Availability", Inches(7.8)),
    ]
    for label, desc, left_pos in milestones:
        dot = slide4.shapes.add_shape(
            MSO_SHAPE.OVAL, left_pos, Inches(3.25), Inches(0.6), Inches(0.6)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor(0xE7, 0x4C, 0x3C)
        dot.line.fill.background()

        lbl = slide4.shapes.add_textbox(left_pos - Inches(0.3), Inches(2.3), Inches(1.5), Inches(0.8))
        tf_lbl = lbl.text_frame
        tf_lbl.word_wrap = True
        p_date = tf_lbl.paragraphs[0]
        p_date.text = label
        p_date.alignment = PP_ALIGN.CENTER
        r_date = p_date.runs[0]
        r_date.font.size = Pt(11)
        r_date.font.bold = True
        r_date.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

        p_desc = tf_lbl.add_paragraph()
        p_desc.text = desc
        p_desc.alignment = PP_ALIGN.CENTER
        r_desc = p_desc.runs[0]
        r_desc.font.size = Pt(10)
        r_desc.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)

    # ====================== SLIDE 5: Contact / Closing ======================
    slide5 = prs.slides.add_slide(blank_layout)
    fill5 = slide5.background.fill
    fill5.solid()
    fill5.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Big "Thank You" text
    ty = slide5.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7.0), Inches(2.0))
    tf_ty = ty.text_frame
    tf_ty.word_wrap = True
    p_ty = tf_ty.paragraphs[0]
    p_ty.text = "Thank You"
    p_ty.alignment = PP_ALIGN.CENTER
    r_ty = p_ty.runs[0]
    r_ty.font.name = "Arial"
    r_ty.font.size = Pt(44)
    r_ty.font.bold = True
    r_ty.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    # Contact info box
    contact = slide5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(4.0), Inches(5.0), Inches(2.0)
    )
    contact.fill.solid()
    contact.fill.fore_color.rgb = RGBColor(0x34, 0x49, 0x5E)
    contact.line.fill.background()
    tf_contact = contact.text_frame
    tf_contact.word_wrap = True
    p_c1 = tf_contact.paragraphs[0]
    p_c1.text = "Sarah Chen | Lead Designer"
    p_c1.alignment = PP_ALIGN.CENTER
    r_c1 = p_c1.runs[0]
    r_c1.font.size = Pt(16)
    r_c1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p_c2 = tf_contact.add_paragraph()
    p_c2.text = "sarah.chen@overlaytech.com"
    p_c2.alignment = PP_ALIGN.CENTER
    r_c2 = p_c2.runs[0]
    r_c2.font.size = Pt(14)
    r_c2.font.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)

    # Decorative diamond
    diamond = slide5.shapes.add_shape(
        MSO_SHAPE.DIAMOND, Inches(8.5), Inches(0.5), Inches(1.0), Inches(1.0)
    )
    diamond.fill.solid()
    diamond.fill.fore_color.rgb = RGBColor(0xF3, 0x9C, 0x12)
    diamond.line.fill.background()

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
