"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 164, I want the heading to be split over two lines so it reads exactly like this:
Line 1: “Title —”
Line 2: “Subtitle”
How do I insert that line break inside the title text box in LibreOffice Impress?
Generated: 2025-09-10 16:56:54
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation

def verify_title_break(file_path: str) -> float:
    """Verify that on slide 164 the heading is split over two lines exactly as:
    Line 1: "Title —"
    Line 2: "Subtitle"

    Progressive scoring (total 1.0):
      • 0.2 – Presentation contains at least 164 slides (so slide 164 exists)
      • 0.3 – Slide 164 contains a text shape that includes BOTH words ‘Title’ and ‘Subtitle’
      • 0.5 – That shape’s text is split into exactly two lines that match the required strings
    """
    score = 0.0
    max_score = 1.0

    # ------------------------------------------------------------
    # Step 1 – Load the presentation (no points for loading itself)
    # ------------------------------------------------------------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error opening presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    total_slides = len(prs.slides)
    print(f"Loaded presentation with {total_slides} slides")

    # ------------------------------------------------------------
    # Step 2 – Ensure slide 164 exists (index 163)
    # ------------------------------------------------------------
    if total_slides >= 164:
        print("✓ Presentation has at least 164 slides (0.2 points)")
        score += 0.2
    else:
        print("✗ Presentation has fewer than 164 slides – cannot verify heading")
        print(f"REWARD: {score}")
        return score  # early exit, no more points possible

    # ------------------------------------------------------------
    # Step 3 – Locate heading text on slide 164
    # ------------------------------------------------------------
    slide = prs.slides[163]  # zero-based index
    candidate_text = None

    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text.strip().replace("\r", "")  # normalise
            if text and ("Title" in text and "Subtitle" in text):
                candidate_text = text
                break  # first matching shape is enough

    if candidate_text is None:
        print("✗ No heading shape containing both ‘Title’ and ‘Subtitle’ found on slide 164")
        print(f"REWARD: {score}")
        return score

    print("✓ Found heading shape on slide 164 (0.3 points)")
    score += 0.3

    # ------------------------------------------------------------
    # Step 4 – Verify exact two-line content with correct dash
    # ------------------------------------------------------------
    lines = [ln.strip() for ln in candidate_text.split("\n") if ln.strip()]
    print(f"Extracted heading lines: {lines}")

    if len(lines) == 2:
        first_ok = lines[0] == "Title —"  # EM dash U+2014
        second_ok = lines[1] == "Subtitle"
        if first_ok and second_ok:
            print("✓ Heading lines match exactly (0.5 points)")
            score += 0.5
        else:
            # Partial credit if only dash style differs but words correct
            dash_normalised = lines[0].replace("—", "-").replace("–", "-")
            if dash_normalised == "Title -" and second_ok:
                print("✓ Heading text correct but dash style differs (0.25 points)")
                score += 0.25  # partial credit
            else:
                print("✗ Heading lines do not match exactly – no additional points")
    else:
        print("✗ Heading is not split into exactly two non-empty lines – no additional points")

    # ------------------------------------------------------------
    # Final score (capped at 1.0)
    # ------------------------------------------------------------
    final_score = min(score, max_score)
    print(f"REWARD: {final_score}")
    return final_score

# ------------------------------------------------------------
# Execution entry point
# ------------------------------------------------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/on_slide_164_i_want_the_heading_to_be_split_over_two_lines_so_it_reads_exactly_like_this_line_1_titl_golden.pptx"
    verify_title_break(FILE_PATH)

