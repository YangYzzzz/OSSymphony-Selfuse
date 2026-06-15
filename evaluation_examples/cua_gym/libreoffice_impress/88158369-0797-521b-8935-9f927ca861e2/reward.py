"""
Reward Script: Interactive presentation with action buttons
Task ID: impress_gf5_031
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Master slide has Previous, Next, Home buttons with correct text
  Component 2 (0.25): Master buttons have correct action settings (ppaction links)
  Component 3 (0.30): Glossary button present on exactly slides 3, 7, 10, 15, 18
  Component 4 (0.20): Glossary buttons link to /home/user/glossary.pdf
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_031'


def persist_app_state(domain):
    """Save any unsaved LibreOffice changes before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.oxml.ns import qn
        import lxml.etree as etree
    except ImportError as e:
        print(f"CRITICAL: Missing library: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: 20 slides
    if len(prs.slides) != 20:
        print(f"PRECONDITION FAIL: Expected 20 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # --- Component 1: Master slide has Previous, Next, Home buttons (0.25 points) ---
    try:
        master = prs.slide_masters[0]
        master_button_texts = {}
        for shape in master.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip().lower()
                if text in ('previous', 'next', 'home'):
                    master_button_texts[text] = shape

        found_buttons = set(master_button_texts.keys())
        expected_buttons = {'previous', 'next', 'home'}

        if found_buttons == expected_buttons:
            print(f"PASS: Component 1 — All 3 master buttons found: {found_buttons} (0.25 pts)")
            total_score += 0.25
        elif len(found_buttons) > 0:
            partial = 0.25 * len(found_buttons & expected_buttons) / 3.0
            print(f"PARTIAL: Component 1 — Found {found_buttons}, missing {expected_buttons - found_buttons} ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — No navigation buttons found on master slide")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # --- Component 2: Master buttons have correct action settings (0.25 points) ---
    try:
        action_checks = {
            'previous': 'previousslide',
            'next': 'nextslide',
            'home': 'firstslide',
        }
        actions_correct = 0
        actions_total = 0

        for btn_text, expected_action in action_checks.items():
            if btn_text not in master_button_texts:
                continue
            actions_total += 1
            shape = master_button_texts[btn_text]
            el = shape._element
            xml_str = etree.tostring(el).decode()

            if expected_action in xml_str.lower():
                actions_correct += 1
                print(f"  PASS: '{btn_text}' button action contains '{expected_action}'")
            else:
                print(f"  FAIL: '{btn_text}' button missing action '{expected_action}'")

        if actions_total > 0 and actions_correct == actions_total:
            print(f"PASS: Component 2 — All {actions_correct} button actions correct (0.25 pts)")
            total_score += 0.25
        elif actions_correct > 0:
            partial = 0.25 * actions_correct / 3.0
            print(f"PARTIAL: Component 2 — {actions_correct}/{actions_total} actions correct ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No correct button actions found")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # --- Component 3: Glossary button on exactly slides 3, 7, 10, 15, 18 (0.30 points) ---
    try:
        expected_glossary_slides = {3, 7, 10, 15, 18}
        actual_glossary_slides = set()

        for idx in range(len(prs.slides)):
            slide = prs.slides[idx]
            slide_num = idx + 1
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip().lower()
                    if text == 'glossary' or 'glossary' in shape.name.lower():
                        actual_glossary_slides.add(slide_num)

        correct_placements = expected_glossary_slides & actual_glossary_slides
        wrong_placements = actual_glossary_slides - expected_glossary_slides

        if actual_glossary_slides == expected_glossary_slides:
            print(f"PASS: Component 3 — Glossary buttons on exactly slides {sorted(expected_glossary_slides)} (0.30 pts)")
            total_score += 0.30
        elif len(correct_placements) > 0 and len(wrong_placements) == 0:
            partial = 0.30 * len(correct_placements) / len(expected_glossary_slides)
            print(f"PARTIAL: Component 3 — Glossary on {sorted(correct_placements)}, missing {sorted(expected_glossary_slides - correct_placements)} ({partial:.3f} pts)")
            total_score += partial
        elif len(correct_placements) > 0:
            # Some correct but also some wrong placements — reduce credit
            partial = 0.30 * len(correct_placements) / len(expected_glossary_slides) * 0.5
            print(f"PARTIAL: Component 3 — Glossary on {sorted(actual_glossary_slides)}, expected {sorted(expected_glossary_slides)} ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No glossary buttons found on expected slides. Found on: {sorted(actual_glossary_slides)}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # --- Component 4: Glossary buttons link to /home/user/glossary.pdf (0.20 points) ---
    try:
        links_correct = 0
        links_checked = 0

        for idx in [s - 1 for s in expected_glossary_slides]:
            if idx >= len(prs.slides):
                continue
            slide = prs.slides[idx]
            for shape in slide.shapes:
                is_glossary = False
                if shape.has_text_frame and shape.text_frame.text.strip().lower() == 'glossary':
                    is_glossary = True
                if 'glossary' in shape.name.lower():
                    is_glossary = True

                if is_glossary:
                    links_checked += 1
                    # Check hyperlink via XML
                    el = shape._element
                    xml_str = etree.tostring(el).decode()

                    # Find hlinkClick rId, then resolve via slide rels
                    hlink_els = el.findall('.//' + qn('a:hlinkClick'))
                    link_found = False
                    for hlink in hlink_els:
                        r_id = hlink.get(qn('r:id'))
                        if r_id:
                            # Look up relationship target
                            try:
                                rel = slide.part.rels[r_id]
                                target = str(rel.target_ref)
                                if 'glossary.pdf' in target:
                                    links_correct += 1
                                    link_found = True
                                    print(f"  PASS: Slide {idx+1} glossary links to {target}")
                                else:
                                    print(f"  FAIL: Slide {idx+1} glossary links to {target}, expected glossary.pdf")
                            except KeyError:
                                print(f"  FAIL: Slide {idx+1} glossary has rId={r_id} but no matching rel")
                        # Also check for action-based hyperlink
                        action = hlink.get('action', '')
                        if 'glossary.pdf' in action:
                            links_correct += 1
                            link_found = True
                            print(f"  PASS: Slide {idx+1} glossary action links to glossary.pdf")
                    if not link_found and not hlink_els:
                        print(f"  FAIL: Slide {idx+1} glossary has no hlinkClick")

        if links_checked > 0 and links_correct == links_checked:
            print(f"PASS: Component 4 — All {links_correct} glossary buttons link to glossary.pdf (0.20 pts)")
            total_score += 0.20
        elif links_correct > 0:
            partial = 0.20 * links_correct / max(links_checked, len(expected_glossary_slides))
            print(f"PARTIAL: Component 4 — {links_correct}/{links_checked} glossary links correct ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — No correct glossary hyperlinks found (checked {links_checked} buttons)")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
