"""
Initial Setup: Create executive_summary.pptx with 20 slides of business content.
Task ID: impress_gf5_031
Domain: libreoffice_impress
No navigation buttons. Just content slides open in LibreOffice Impress.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_031'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                alignment=PP_ALIGN.LEFT, color=None):
    """Helper to add a textbox with a single run."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=14):
    """Add a bulleted list textbox."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(font_size)
    return txBox

def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide content data - 20 slides of executive summary content
    slides_data = [
        {
            "title": "Q4 2025 Executive Summary",
            "subtitle": "Global Operations & Strategic Initiatives Review",
            "type": "title"
        },
        {
            "title": "Agenda",
            "bullets": [
                "Financial Performance Overview",
                "Regional Market Analysis",
                "Product Portfolio Update",
                "Strategic Partnership Review",
                "Workforce & Talent Development",
                "Technology Infrastructure Roadmap",
                "Risk Assessment & Mitigation",
                "Q1 2026 Strategic Priorities"
            ],
            "type": "content"
        },
        {
            "title": "Financial Performance Overview",
            "bullets": [
                "Total Revenue: $487.3M (up 12.4% YoY)",
                "Gross Margin: 42.8% (improved from 39.1%)",
                "Operating Expenses: $198.2M (within budget)",
                "EBITDA: $124.6M (exceeding target by 8%)",
                "Free Cash Flow: $67.9M",
                "Customer Acquisition Cost: $142 (down 15%)"
            ],
            "type": "content"
        },
        {
            "title": "Revenue Breakdown by Division",
            "bullets": [
                "Enterprise Solutions: $201.4M (41.3%)",
                "Cloud Services: $156.8M (32.2%)",
                "Professional Services: $78.3M (16.1%)",
                "Managed Infrastructure: $50.8M (10.4%)"
            ],
            "type": "content"
        },
        {
            "title": "North America Market Analysis",
            "bullets": [
                "Market share increased to 18.7% from 16.2%",
                "New enterprise clients: 47 (target was 40)",
                "Customer retention rate: 94.3%",
                "Average deal size: $1.2M (up from $980K)",
                "Sales pipeline: $890M qualified opportunities"
            ],
            "type": "content"
        },
        {
            "title": "EMEA Market Performance",
            "bullets": [
                "Revenue: $112.6M (23.1% of total)",
                "Germany and UK remain top markets",
                "Opened new office in Stockholm",
                "Regulatory compliance: GDPR audit passed",
                "Hired 34 sales representatives across 6 countries"
            ],
            "type": "content"
        },
        {
            "title": "Asia-Pacific Expansion Update",
            "bullets": [
                "Revenue growth: 28.4% YoY (fastest region)",
                "Singapore hub fully operational",
                "Japan partnership with NTT Data signed",
                "Australia: 12 new government contracts",
                "India development center: 200+ engineers"
            ],
            "type": "content"
        },
        {
            "title": "Product Portfolio: Enterprise Suite 5.0",
            "bullets": [
                "General availability: January 15, 2026",
                "142 new features shipped in Q4",
                "AI-powered analytics module completed",
                "Performance improvements: 3x faster query execution",
                "99.97% uptime SLA maintained",
                "Security certifications: SOC 2 Type II, ISO 27001"
            ],
            "type": "content"
        },
        {
            "title": "Cloud Migration Progress",
            "bullets": [
                "78% of enterprise clients migrated to cloud",
                "Hybrid deployment option launched",
                "Average migration time: 6.2 weeks (down from 11)",
                "Zero data loss incidents during migration",
                "Cloud infrastructure costs optimized by 22%"
            ],
            "type": "content"
        },
        {
            "title": "Strategic Partnerships",
            "bullets": [
                "Microsoft Azure: Platinum Partner status achieved",
                "AWS Marketplace: Revenue up 45%",
                "Salesforce integration: 2,400 joint customers",
                "SAP partnership: Co-development agreement signed",
                "Deloitte: Preferred implementation partner"
            ],
            "type": "content"
        },
        {
            "title": "Customer Success Highlights",
            "bullets": [
                "Net Promoter Score: 72 (industry avg: 41)",
                "Customer Health Score: 87/100",
                "Support ticket resolution: 4.2 hours average",
                "Quarterly Business Reviews: 98% completion",
                "Case Study: Goldman Sachs saved $12M annually",
                "Case Study: Siemens reduced cycle time by 40%"
            ],
            "type": "content"
        },
        {
            "title": "Workforce & Talent Overview",
            "bullets": [
                "Total headcount: 3,847 (up 312 from Q3)",
                "Engineering team: 1,420 (37%)",
                "Voluntary turnover: 8.2% (industry avg: 13%)",
                "Diversity: 44% female leadership (up from 38%)",
                "Employee engagement score: 4.3/5.0",
                "Training hours per employee: 42 per quarter"
            ],
            "type": "content"
        },
        {
            "title": "Technology Infrastructure Roadmap",
            "bullets": [
                "Kubernetes migration: 92% complete",
                "Edge computing pilot in 3 regions",
                "Zero-trust security framework deployed",
                "API gateway: 2.8B requests per day",
                "Data lake consolidation: Phase 2 complete",
                "Disaster recovery: RPO < 15 minutes achieved"
            ],
            "type": "content"
        },
        {
            "title": "Research & Development Investments",
            "bullets": [
                "R&D spending: $73.4M (15.1% of revenue)",
                "AI/ML team expanded to 85 researchers",
                "12 patent applications filed in Q4",
                "University partnerships: MIT, Stanford, ETH Zurich",
                "Open-source contributions: 340 PRs merged"
            ],
            "type": "content"
        },
        {
            "title": "Risk Assessment Summary",
            "bullets": [
                "Cybersecurity: 0 critical incidents (target met)",
                "Supply chain: Component availability improved",
                "Regulatory: EU AI Act compliance plan approved",
                "Currency exposure: Hedged 80% of EUR/GBP",
                "Talent: Key person risk mitigated with succession plans",
                "Competition: Market position strengthened"
            ],
            "type": "content"
        },
        {
            "title": "Sustainability & ESG Progress",
            "bullets": [
                "Carbon emissions: Reduced 18% from baseline",
                "100% renewable energy in North America offices",
                "Sustainable procurement: 72% of vendors certified",
                "Community investment: $4.2M in education programs",
                "ESG rating upgraded to AA by MSCI"
            ],
            "type": "content"
        },
        {
            "title": "Competitive Landscape",
            "bullets": [
                "Won 23 competitive displacements in Q4",
                "Key wins against Oracle (8), ServiceNow (6), Workday (5)",
                "Gartner Magic Quadrant: Leader position maintained",
                "Forrester Wave: Strong Performer → Leader",
                "G2 Crowd: Highest satisfaction score in category"
            ],
            "type": "content"
        },
        {
            "title": "Q1 2026 Strategic Priorities",
            "bullets": [
                "Launch Enterprise Suite 5.0 globally",
                "Expand APAC sales team by 50 reps",
                "Complete remaining cloud migrations",
                "Achieve FedRAMP High authorization",
                "Integrate AI copilot into all product lines",
                "Close Series F funding round ($200M target)"
            ],
            "type": "content"
        },
        {
            "title": "Key Performance Indicators — Q1 Targets",
            "bullets": [
                "Revenue: $510M (quarterly target)",
                "New logo acquisition: 55 enterprise clients",
                "Cloud revenue share: 38% of total",
                "Employee NPS: > 75",
                "Customer churn: < 3%",
                "Product release velocity: 150+ features"
            ],
            "type": "content"
        },
        {
            "title": "Thank You — Questions & Discussion",
            "subtitle": "Contact: strategy@acmecorp.com | Investor Relations: ir@acmecorp.com",
            "type": "title"
        },
    ]

    for idx, sdata in enumerate(slides_data):
        if sdata["type"] == "title":
            slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
            slide.shapes.title.text = sdata["title"]
            if "subtitle" in sdata:
                slide.placeholders[1].text = sdata["subtitle"]
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
            # Title
            add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                        sdata["title"], font_size=28, bold=True,
                        color=RGBColor(0x1A, 0x3C, 0x6E))
            # Divider line
            line = slide.shapes.add_shape(
                1,  # MSO_SHAPE.RECTANGLE
                Inches(0.8), Inches(1.2), Inches(11), Inches(0.03)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(0x2E, 0x75, 0xB6)
            line.line.fill.background()

            # Bullet content
            add_bullet_list(slide, Inches(1.0), Inches(1.5), Inches(10), Inches(5.0),
                            sdata["bullets"], font_size=16)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Also create glossary.pdf placeholder (the task references it)
    try:
        from fpdf import FPDF
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font('Helvetica', 'B', 24)
        pdf.cell(0, 20, 'Glossary of Terms', ln=True, align='C')
        pdf.set_font('Helvetica', '', 12)
        terms = [
            ("EBITDA", "Earnings Before Interest, Taxes, Depreciation, and Amortization"),
            ("YoY", "Year over Year comparison"),
            ("SLA", "Service Level Agreement"),
            ("RPO", "Recovery Point Objective"),
            ("NPS", "Net Promoter Score"),
            ("FedRAMP", "Federal Risk and Authorization Management Program"),
            ("ESG", "Environmental, Social, and Governance"),
            ("SOC 2", "System and Organization Controls 2"),
            ("APAC", "Asia-Pacific region"),
            ("EMEA", "Europe, Middle East, and Africa"),
        ]
        for term, definition in terms:
            pdf.set_font('Helvetica', 'B', 12)
            pdf.cell(0, 10, f'{term}:', ln=False)
            pdf.set_font('Helvetica', '', 12)
            pdf.cell(0, 10, f'  {definition}', ln=True)
        pdf.output(f'{WORKDIR}/glossary.pdf')
        print(f'Glossary PDF created: {WORKDIR}/glossary.pdf')
    except Exception as e:
        print(f'Warning: Could not create glossary.pdf: {e}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')

create_initial()
