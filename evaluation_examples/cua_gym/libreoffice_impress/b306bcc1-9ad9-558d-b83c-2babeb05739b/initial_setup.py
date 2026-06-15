"""
Initial Setup: Create a 6-slide sales review presentation
Task ID: impstruct_011
Domain: libreoffice_impress
Slides: Sales Review 2025, Q1, Q2, Q3, Q4, Summary
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impstruct_011'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_content_slide(prs, layout_idx, title_text, body_lines):
    """Add a slide with title and bullet-point body content."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Sales Review 2025"
    slide1.placeholders[1].text = "Annual Performance Overview\nPrepared by the Revenue Analytics Team"

    # --- Slide 2: Q1 ---
    add_content_slide(prs, 1, "Q1", [
        "Total Revenue: $2.34M (+8% YoY)",
        "New Accounts: 47 enterprise clients signed",
        "Top Region: North America ($1.12M)",
        "Key Win: Meridian Healthcare 3-year contract ($420K)",
        "Churn Rate: 3.1% (down from 4.2%)",
        "Sales Team Headcount: 28 reps across 4 regions",
    ])

    # --- Slide 3: Q2 ---
    add_content_slide(prs, 1, "Q2", [
        "Total Revenue: $2.78M (+12% YoY)",
        "New Accounts: 53 enterprise clients signed",
        "Top Region: EMEA ($985K, strongest quarter)",
        "Key Win: Cascade Financial annual deal ($310K)",
        "Product Launch: Analytics Pro drove 22% of new pipeline",
        "Average Deal Size: $58.7K (up from $51.2K in Q1)",
    ])

    # --- Slide 4: Q3 ---
    add_content_slide(prs, 1, "Q3", [
        "Total Revenue: $3.15M (+15% YoY)",
        "New Accounts: 61 enterprise clients signed",
        "Top Region: APAC ($1.05M, first time leading)",
        "Key Win: Thornton Industries multi-year deal ($575K)",
        "Upsell Revenue: $890K from existing customer expansions",
        "Pipeline Coverage: 3.4x for Q4 target",
    ])

    # --- Slide 5: Q4 ---
    add_content_slide(prs, 1, "Q4", [
        "Total Revenue: $3.62M (+18% YoY)",
        "New Accounts: 72 enterprise clients signed",
        "Top Region: North America ($1.48M, holiday surge)",
        "Key Win: Vanguard Logistics global rollout ($680K)",
        "Annual Recurring Revenue: $11.89M total",
        "Quota Attainment: 112% of annual target",
    ])

    # --- Slide 6: Summary ---
    add_content_slide(prs, 1, "Summary", [
        "Full-Year Revenue: $11.89M (+13.2% YoY)",
        "Total New Accounts: 233 enterprise clients",
        "Net Revenue Retention: 118%",
        "Top Performing Region: North America ($4.43M)",
        "Fastest Growing Region: APAC (+34% YoY)",
        "2026 Target: $14.5M with 15% headcount growth",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
