"""
Reward Script: Swap slide positions in a presentation
Task ID: impstruct_011
Domain: libreoffice_impress
Scoring:
  - Component 1: Slide 2 title is 'Q3' (0.25 pts)
  - Component 2: Slide 3 title is 'Q4' (0.25 pts)
  - Component 3: Slide 4 title is 'Q1' (0.25 pts)
  - Component 4: Slide 5 title is 'Q2' (0.25 pts)
  Preconditions: 6 slides, slide 1 = 'Sales Review 2025', slide 6 = 'Summary'
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impstruct_011'


def persist_app_state(domain: str):
    """Attempt to save any unsaved LibreOffice edits via Ctrl+S."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_slide_title(slide):
    """Extract the first non-empty text from a slide's shapes (title)."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    return text
    return ""


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Must have exactly 6 slides
    num_slides = len(prs.slides)
    if num_slides != 6:
        print(f"PRECONDITION FAIL: Expected 6 slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)
    titles = [get_slide_title(s) for s in slides]
    print(f"Slide titles found: {titles}")

    # Precondition: Slide 1 must be 'Sales Review 2025' (unchanged)
    if titles[0] != 'Sales Review 2025':
        print(f"PRECONDITION FAIL: Slide 1 title expected 'Sales Review 2025', found '{titles[0]}'")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Slide 6 must be 'Summary' (unchanged)
    if titles[5] != 'Summary':
        print(f"PRECONDITION FAIL: Slide 6 title expected 'Summary', found '{titles[5]}'")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide 2 should now be 'Q3' (was at position 4, swapped with position 2)
    # In initial_env, slide 2 is 'Q1' -- so this check FAILS on initial, PASSES on golden
    try:
        if titles[1] == 'Q3':
            print(f"PASS: Component 1 -- Slide 2 title is 'Q3' (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 -- Slide 2 title expected 'Q3', found '{titles[1]}'")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Slide 3 should now be 'Q4' (was at position 5, swapped with position 3)
    # In initial_env, slide 3 is 'Q2' -- so this check FAILS on initial, PASSES on golden
    try:
        if titles[2] == 'Q4':
            print(f"PASS: Component 2 -- Slide 3 title is 'Q4' (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 -- Slide 3 title expected 'Q4', found '{titles[2]}'")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Slide 4 should now be 'Q1' (was at position 2, swapped with position 4)
    # In initial_env, slide 4 is 'Q3' -- so this check FAILS on initial, PASSES on golden
    try:
        if titles[3] == 'Q1':
            print(f"PASS: Component 3 -- Slide 4 title is 'Q1' (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 -- Slide 4 title expected 'Q1', found '{titles[3]}'")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Slide 5 should now be 'Q2' (was at position 3, swapped with position 5)
    # In initial_env, slide 5 is 'Q4' -- so this check FAILS on initial, PASSES on golden
    try:
        if titles[4] == 'Q2':
            print(f"PASS: Component 4 -- Slide 5 title is 'Q2' (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 4 -- Slide 5 title expected 'Q2', found '{titles[4]}'")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist any unsaved GUI state before verification
persist_app_state("libreoffice_impress")

# Run verification
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
