"""
Reward Script: Create a formatted comparison table on slide 3
Task ID: impress_gf2_006
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Table exists on slide 3 with 7 rows x 5 columns
  Component 2 (0.25): Header row formatting (dark bg, white bold text, correct labels)
  Component 3 (0.25): Data row formatting (alternating colors, center-aligned)
  Component 4 (0.25): Correct checkmark/cross symbols in data cells
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_006'


def persist_app_state(domain):
    """Save any unsaved LibreOffice changes before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def find_table_on_slide(slide):
    """Find the first table shape on a slide, or None."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return shape.table
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]  # 0-indexed, slide 3

    # Component 1: Table exists on slide 3 with correct dimensions (0.25 points)
    try:
        table = find_table_on_slide(slide3)
        if table is None:
            print("FAIL: Component 1 — No table found on slide 3")
        else:
            num_rows = len(table.rows)
            num_cols = len(table.columns)
            if num_rows == 7 and num_cols == 5:
                print(f"PASS: Component 1 — Table found with {num_rows} rows x {num_cols} cols (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 — Table has {num_rows} rows x {num_cols} cols, expected 7x5")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Gate: if no table, remaining checks cannot proceed
    if table is None:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: Header row formatting (0.25 points)
    # Dark background #1E3A5F, white bold text, correct column names
    try:
        expected_headers = ['Feature', 'Starter', 'Professional', 'Enterprise', 'Custom']
        header_checks_passed = 0
        total_header_checks = len(expected_headers) * 3  # text + bg + font per cell

        for c in range(min(len(table.columns), 5)):
            cell = table.cell(0, c)

            # Check header text
            if cell.text.strip() == expected_headers[c]:
                header_checks_passed += 1

            # Check background color (#1E3A5F)
            try:
                if cell.fill.type is not None and str(cell.fill.fore_color.rgb) == '1E3A5F':
                    header_checks_passed += 1
            except Exception:
                pass

            # Check white bold text
            try:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        is_white = False
                        is_bold = False
                        try:
                            if str(run.font.color.rgb) == 'FFFFFF':
                                is_white = True
                        except Exception:
                            pass
                        if run.font.bold is True:
                            is_bold = True
                        if is_white and is_bold:
                            header_checks_passed += 1
                            break
                    break  # only check first paragraph
            except Exception:
                pass

        # Award proportional credit
        header_ratio = header_checks_passed / total_header_checks if total_header_checks > 0 else 0
        header_score = 0.25 * header_ratio
        if header_ratio >= 0.8:
            # Round up if mostly correct
            header_score = 0.25
        total_score += header_score
        print(f"{'PASS' if header_ratio >= 0.8 else 'PARTIAL'}: Component 2 — Header formatting {header_checks_passed}/{total_header_checks} checks ({header_score:.3f} pts)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Data row formatting — alternating colors and center alignment (0.25 points)
    try:
        # Expected: odd data rows (1,3,5) = #DBEAFE, even data rows (2,4,6) = #FFFFFF
        expected_colors = {1: 'DBEAFE', 2: 'FFFFFF', 3: 'DBEAFE', 4: 'FFFFFF', 5: 'DBEAFE', 6: 'FFFFFF'}
        format_checks_passed = 0
        total_format_checks = 0

        for r in range(1, min(len(table.rows), 7)):
            # Check background color of first cell in each data row
            total_format_checks += 1
            try:
                cell = table.cell(r, 0)
                if cell.fill.type is not None and str(cell.fill.fore_color.rgb) == expected_colors[r]:
                    format_checks_passed += 1
                else:
                    try:
                        actual = str(cell.fill.fore_color.rgb) if cell.fill.type is not None else 'None'
                    except Exception:
                        actual = 'unknown'
                    print(f"  Row {r} bg: expected {expected_colors[r]}, got {actual}")
            except Exception:
                pass

            # Check center alignment
            total_format_checks += 1
            try:
                cell = table.cell(r, 0)
                for para in cell.text_frame.paragraphs:
                    if para.alignment == PP_ALIGN.CENTER:
                        format_checks_passed += 1
                    break
            except Exception:
                pass

        format_ratio = format_checks_passed / total_format_checks if total_format_checks > 0 else 0
        format_score = 0.25 * format_ratio
        if format_ratio >= 0.8:
            format_score = 0.25
        total_score += format_score
        print(f"{'PASS' if format_ratio >= 0.8 else 'PARTIAL'}: Component 3 — Data row formatting {format_checks_passed}/{total_format_checks} checks ({format_score:.3f} pts)")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Correct checkmark/cross content in data cells (0.25 points)
    try:
        expected_data = [
            ['Cloud Storage (50 GB)', '\u2713', '\u2713', '\u2713', '\u2713'],
            ['Real-Time Collaboration', '\u2717', '\u2713', '\u2713', '\u2713'],
            ['Advanced Analytics Dashboard', '\u2717', '\u2717', '\u2713', '\u2713'],
            ['Single Sign-On (SSO)', '\u2717', '\u2717', '\u2713', '\u2713'],
            ['Dedicated Account Manager', '\u2717', '\u2717', '\u2717', '\u2713'],
            ['Custom API Integrations', '\u2717', '\u2717', '\u2717', '\u2713'],
        ]
        content_checks_passed = 0
        total_content_checks = 0

        for r_idx, row_data in enumerate(expected_data):
            r = r_idx + 1  # table row index (skip header)
            if r >= len(table.rows):
                break
            for c_idx, expected_text in enumerate(row_data):
                if c_idx >= len(table.columns):
                    break
                total_content_checks += 1
                actual_text = table.cell(r, c_idx).text.strip()
                if c_idx == 0:
                    # Feature name — check containment (flexible)
                    if expected_text.lower() in actual_text.lower() or actual_text.lower() in expected_text.lower():
                        content_checks_passed += 1
                    else:
                        print(f"  Cell({r},{c_idx}): expected '{expected_text}', got '{actual_text}'")
                else:
                    # Checkmark/cross — check the symbol
                    if expected_text in actual_text:
                        content_checks_passed += 1
                    else:
                        print(f"  Cell({r},{c_idx}): expected '{expected_text}', got '{actual_text}'")

        content_ratio = content_checks_passed / total_content_checks if total_content_checks > 0 else 0
        content_score = 0.25 * content_ratio
        if content_ratio >= 0.8:
            content_score = 0.25
        total_score += content_score
        print(f"{'PASS' if content_ratio >= 0.8 else 'PARTIAL'}: Component 4 — Data content {content_checks_passed}/{total_content_checks} checks ({content_score:.3f} pts)")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state('libreoffice_impress')

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
