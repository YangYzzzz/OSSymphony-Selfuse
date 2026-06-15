"""
Initial Setup: Create a 6-slide Pricing Plans presentation with empty Slide 3 content area.
Task ID: impress_gf2_006
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_006'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_text(slide, title_text, subtitle_text=None):
    """Set title and optional subtitle on a slide with placeholders."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    if subtitle_text and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle_text


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    add_title_text(slide1, "CloudSync Pro", "Pricing Plans & Feature Overview")
    # Dark background
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
    # Make title text white
    for run in slide1.shapes.title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.size = Pt(40)
    if len(slide1.placeholders) > 1:
        for run in slide1.placeholders[1].text_frame.paragraphs[0].runs:
            run.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
            run.font.size = Pt(22)

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_textbox(slide2, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                "Why Choose CloudSync Pro?", font_size=32, bold=True,
                color=RGBColor(0x1E, 0x3A, 0x5F))
    bullet_items = [
        "Trusted by over 12,000 businesses worldwide since 2019",
        "99.99% uptime SLA with global data center redundancy",
        "SOC 2 Type II and ISO 27001 certified infrastructure",
        "Dedicated customer success managers for Enterprise plans",
        "Seamless migration from Dropbox, Google Drive, and OneDrive",
    ]
    y = Inches(1.8)
    for item in bullet_items:
        add_textbox(slide2, Inches(1.2), y, Inches(10), Inches(0.5),
                    f"  {item}", font_size=16,
                    color=RGBColor(0x33, 0x33, 0x33))
        y += Inches(0.55)

    # --- Slide 3: Feature Comparison (EMPTY - no table!) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_textbox(slide3, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
                "Feature Comparison", font_size=32, bold=True,
                color=RGBColor(0x1E, 0x3A, 0x5F))
    # Empty content area - agent must create the table here

    # --- Slide 4: Pricing Tiers ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_textbox(slide4, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                "Pricing Tiers", font_size=32, bold=True,
                color=RGBColor(0x1E, 0x3A, 0x5F))
    tiers = [
        ("Starter", "$9.99/mo", "Perfect for individuals and freelancers"),
        ("Professional", "$29.99/mo", "Ideal for small teams up to 25 users"),
        ("Enterprise", "$79.99/mo", "Built for organizations with 100+ users"),
        ("Custom", "Contact Us", "Tailored solutions for unique requirements"),
    ]
    y = Inches(1.8)
    for name, price, desc in tiers:
        add_textbox(slide4, Inches(1.0), y, Inches(3), Inches(0.5),
                    name, font_size=22, bold=True,
                    color=RGBColor(0x1E, 0x3A, 0x5F))
        add_textbox(slide4, Inches(4.2), y, Inches(2), Inches(0.5),
                    price, font_size=20, bold=True,
                    color=RGBColor(0x2E, 0x7D, 0x32))
        add_textbox(slide4, Inches(6.5), y + Inches(0.05), Inches(5.5), Inches(0.5),
                    desc, font_size=16,
                    color=RGBColor(0x55, 0x55, 0x55))
        y += Inches(0.7)

    # --- Slide 5: Customer Testimonials ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_textbox(slide5, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                "What Our Customers Say", font_size=32, bold=True,
                color=RGBColor(0x1E, 0x3A, 0x5F))
    testimonials = [
        ('"CloudSync Pro reduced our file management overhead by 40%. The Enterprise plan\'s audit trail saved us during our last compliance review."',
         "- Rachel Kim, CTO at Meridian Analytics"),
        ('"We migrated 2.3 TB of data from Google Drive in under 6 hours. The migration assistant handled everything seamlessly."',
         "- David Okonkwo, IT Director at Bright Horizon Education"),
    ]
    y = Inches(2.0)
    for quote, author in testimonials:
        add_textbox(slide5, Inches(1.2), y, Inches(10), Inches(1.0),
                    quote, font_size=15,
                    color=RGBColor(0x44, 0x44, 0x44))
        add_textbox(slide5, Inches(1.5), y + Inches(0.9), Inches(8), Inches(0.4),
                    author, font_size=14, bold=True,
                    color=RGBColor(0x1E, 0x3A, 0x5F))
        y += Inches(1.8)

    # --- Slide 6: Contact & Next Steps ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    fill6 = slide6.background.fill
    fill6.solid()
    fill6.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
    add_textbox(slide6, Inches(2), Inches(2), Inches(9), Inches(1),
                "Ready to Get Started?", font_size=36, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_textbox(slide6, Inches(2), Inches(3.2), Inches(9), Inches(0.6),
                "Contact our sales team: sales@cloudsyncpro.com", font_size=20,
                color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)
    add_textbox(slide6, Inches(2), Inches(4.0), Inches(9), Inches(0.6),
                "Start your 14-day free trial at cloudsyncpro.com/trial", font_size=18,
                color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
