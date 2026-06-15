"""
Reward Script: Update slide 4's title to 'Product Differentiation' with brand blue color,
and slide 6's title to 'Go-To-Market Strategy' with brand red color.
Task ID: osworld_impress_title_color_match_006
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 4 title text == 'Product Differentiation'     (0.25 pts)
  Component 2: Slide 4 title color == #0057B7 (brand blue)         (0.25 pts)
  Component 3: Slide 6 title text == 'Go-To-Market Strategy'       (0.25 pts)
  Component 4: Slide 6 title color == #D50032 (brand red)          (0.25 pts)
  Total: 1.0
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_color_match_006'

# Expected values from task context
SLIDE4_EXPECTED_TITLE = 'Product Differentiation'
SLIDE4_EXPECTED_COLOR = '0057B7'  # brand blue, matching slide 2 title

SLIDE6_EXPECTED_TITLE = 'Go-To-Market Strategy'
SLIDE6_EXPECTED_COLOR = 'D50032'  # brand red, matching slide 3 title


def get_title_shape(slide):
    """Return the title shape from a slide, or None if not found."""
    for shape in slide.shapes:
        if shape.has_text_frame and 'Title' in shape.name:
            return shape
    return None


def get_shape_title_text(slide):
    """Return the full text of the title shape on a slide."""
    shape = get_title_shape(slide)
    if shape is None:
        return None
    return shape.text_frame.text.strip()


def get_shape_title_color(slide):
    """
    Return the RGB hex string of the title shape's first non-empty run color.
    Returns None if no explicit RGB color is set.
    """
    shape = get_title_shape(slide)
    if shape is None:
        return None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.text.strip():
                try:
                    if run.font.color.type is not None:
                        return str(run.font.color.rgb).upper()
                except Exception:
                    pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition gate: file must be loadable
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition gate: must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"CRITICAL: Expected at least 6 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide4 = prs.slides[3]  # 0-indexed: slide 4
    slide6 = prs.slides[5]  # 0-indexed: slide 6

    # Component 1: Slide 4 title text is 'Product Differentiation' (0.25 pts)
    try:
        actual_title4 = get_shape_title_text(slide4)
        if actual_title4 == SLIDE4_EXPECTED_TITLE:
            print(f"PASS: Component 1 — Slide 4 title text is '{SLIDE4_EXPECTED_TITLE}' (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — Slide 4 title expected '{SLIDE4_EXPECTED_TITLE}', found '{actual_title4}'")
    except Exception as e:
        print(f"ERROR: Component 1 — Could not check slide 4 title text: {e}")

    # Component 2: Slide 4 title color is #0057B7 (brand blue) (0.25 pts)
    try:
        actual_color4 = get_shape_title_color(slide4)
        if actual_color4 == SLIDE4_EXPECTED_COLOR:
            print(f"PASS: Component 2 — Slide 4 title color is #{SLIDE4_EXPECTED_COLOR} (brand blue) (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 — Slide 4 title color expected #{SLIDE4_EXPECTED_COLOR}, found #{actual_color4}")
    except Exception as e:
        print(f"ERROR: Component 2 — Could not check slide 4 title color: {e}")

    # Component 3: Slide 6 title text is 'Go-To-Market Strategy' (0.25 pts)
    try:
        actual_title6 = get_shape_title_text(slide6)
        if actual_title6 == SLIDE6_EXPECTED_TITLE:
            print(f"PASS: Component 3 — Slide 6 title text is '{SLIDE6_EXPECTED_TITLE}' (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 — Slide 6 title expected '{SLIDE6_EXPECTED_TITLE}', found '{actual_title6}'")
    except Exception as e:
        print(f"ERROR: Component 3 — Could not check slide 6 title text: {e}")

    # Component 4: Slide 6 title color is #D50032 (brand red) (0.25 pts)
    try:
        actual_color6 = get_shape_title_color(slide6)
        if actual_color6 == SLIDE6_EXPECTED_COLOR:
            print(f"PASS: Component 4 — Slide 6 title color is #{SLIDE6_EXPECTED_COLOR} (brand red) (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 4 — Slide 6 title color expected #{SLIDE6_EXPECTED_COLOR}, found #{actual_color6}")
    except Exception as e:
        print(f"ERROR: Component 4 — Could not check slide 6 title color: {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in a given env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
