"""
Initial Setup: 8-slide product launch deck with slides 4 and 6 having wrong titles/colors
Task ID: osworld_impress_title_color_match_006
Domain: libreoffice_impress

Initial state:
  - Slide 2 title: "Market Opportunity" in brand blue (#0057B7)
  - Slide 3 title: "Competitive Landscape" in brand red (#D50032)
  - Slide 4 title: "Differentiation" in black (NOT the correct "Product Differentiation")
  - Slide 6 title: "Strategy" in black (NOT the correct "Go-To-Market Strategy")
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_color_match_006'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_title_text_and_color(slide, title_text, color_rgb=None):
    """Set the title placeholder text and optionally a font color."""
    title_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape.shape_type == 13:
            continue
        if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
            ph_idx = shape.placeholder_format.idx
            if ph_idx == 0:  # title placeholder
                title_shape = shape
                break
    # fallback: find the largest text box near the top
    if title_shape is None:
        for shape in slide.shapes:
            if shape.has_text_frame:
                title_shape = shape
                break
    if title_shape is None:
        return

    tf = title_shape.text_frame
    # Clear existing paragraphs and set new text
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ''
    if tf.paragraphs:
        p = tf.paragraphs[0]
        # Remove extra runs
        for run in p.runs[1:]:
            run._r.getparent().remove(run._r)
        if p.runs:
            p.runs[0].text = title_text
            if color_rgb is not None:
                p.runs[0].font.color.rgb = color_rgb
            else:
                # Set to black
                p.runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        else:
            run = p.add_run()
            run.text = title_text
            if color_rgb is not None:
                run.font.color.rgb = color_rgb
            else:
                run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    tf.paragraphs[0].text = title_text
    if color_rgb is not None:
        for run in tf.paragraphs[0].runs:
            run.font.color.rgb = color_rgb
    else:
        for run in tf.paragraphs[0].runs:
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)


def add_slide_with_title_content(prs, title_text, body_lines, title_color=None):
    """Add a slide using Title+Content layout (layout index 1)."""
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)

    # Set title
    title_ph = None
    content_ph = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            title_ph = shape
        elif shape.placeholder_format.idx == 1:
            content_ph = shape

    if title_ph is not None:
        title_ph.text = title_text
        if title_ph.text_frame.paragraphs:
            para = title_ph.text_frame.paragraphs[0]
            if para.runs:
                run = para.runs[0]
                if title_color is not None:
                    run.font.color.rgb = title_color
                else:
                    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
                run.font.size = Pt(28)
                run.font.bold = True

    if content_ph is not None and body_lines:
        tf = content_ph.text_frame
        tf.text = body_lines[0]
        for line in body_lines[1:]:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0

    return slide


def add_title_only_slide(prs, title_text, title_color=None, body_text=None):
    """Add a slide with title-only layout."""
    layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(layout)

    from pptx.util import Inches, Pt
    # Add a title text box at the top
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9.0), Inches(1.2)
    )
    tf = title_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    if title_color is not None:
        run.font.color.rgb = title_color
    else:
        run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    if body_text:
        body_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.8), Inches(9.0), Inches(4.5)
        )
        btf = body_box.text_frame
        btf.word_wrap = True
        for i, line in enumerate(body_text):
            if i == 0:
                bp = btf.paragraphs[0]
            else:
                bp = btf.add_paragraph()
            bp.text = line
            bp.level = 0

    return slide


def create_initial():
    prs = Presentation()
    # Use standard widescreen (default is 10" x 7.5")
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    BRAND_BLUE = RGBColor(0x00, 0x57, 0xB7)
    BRAND_RED = RGBColor(0xD5, 0x00, 0x32)
    BLACK = RGBColor(0x00, 0x00, 0x00)

    # ---------- Slide 1: Title Slide ----------
    layout0 = prs.slide_layouts[0]  # Title Slide
    slide1 = prs.slides.add_slide(layout0)
    for ph in slide1.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "NovaTech X1 Pro"
            for run in ph.text_frame.paragraphs[0].runs:
                run.font.size = Pt(36)
                run.font.bold = True
                run.font.color.rgb = BLACK
        elif ph.placeholder_format.idx == 1:
            ph.text = "Product Launch Presentation — Q2 2025"
            for run in ph.text_frame.paragraphs[0].runs:
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # ---------- Slide 2: Market Opportunity — brand blue title ----------
    add_slide_with_title_content(
        prs,
        title_text="Market Opportunity",
        body_lines=[
            "Total addressable market: $4.2B by 2027",
            "Annual growth rate: 18% CAGR",
            "Key segments: Enterprise (42%), SMB (33%), Consumer (25%)",
            "Primary geographies: North America, EMEA, APAC",
            "Underserved pain points in workflow automation",
        ],
        title_color=BRAND_BLUE,
    )

    # ---------- Slide 3: Competitive Landscape — brand red title ----------
    add_slide_with_title_content(
        prs,
        title_text="Competitive Landscape",
        body_lines=[
            "Competitor A: Strong in legacy enterprise, slow innovation",
            "Competitor B: Consumer-focused, lacks enterprise features",
            "Competitor C: Recent entrant with limited ecosystem",
            "NovaTech X1 Pro uniquely addresses all three gaps",
            "Patent-pending FlowSync™ technology provides moat",
        ],
        title_color=BRAND_RED,
    )

    # ---------- Slide 4: Differentiation — BLACK title (WRONG — task asks to rename + recolor) ----------
    add_slide_with_title_content(
        prs,
        title_text="Differentiation",
        body_lines=[
            "FlowSync™ reduces onboarding time by 60%",
            "Unified API layer across 200+ enterprise integrations",
            "AI-assisted workflow builder: no-code + pro-code",
            "99.99% SLA with real-time failover",
            "Best-in-class security: SOC 2 Type II, ISO 27001",
        ],
        title_color=BLACK,
    )

    # ---------- Slide 5: Customer Validation ----------
    add_slide_with_title_content(
        prs,
        title_text="Customer Validation",
        body_lines=[
            "47 design partners across 12 industry verticals",
            "NPS score: 72 (industry avg: 34)",
            "Average deal size: $85,000 ARR",
            "3 Fortune 500 LOIs signed in Q1 2025",
            "Pilot-to-paid conversion rate: 68%",
        ],
        title_color=BLACK,
    )

    # ---------- Slide 6: Strategy — BLACK title (WRONG — task asks to rename + recolor) ----------
    add_slide_with_title_content(
        prs,
        title_text="Strategy",
        body_lines=[
            "Phase 1: Direct sales to 50 enterprise accounts (Q2–Q3)",
            "Phase 2: Channel partner program launch (Q4 2025)",
            "Phase 3: Self-serve SMB tier via marketplace (Q1 2026)",
            "Target: $12M ARR by end of FY2025",
            "Marketing: thought leadership, analyst relations, events",
        ],
        title_color=BLACK,
    )

    # ---------- Slide 7: Financial Projections ----------
    add_slide_with_title_content(
        prs,
        title_text="Financial Projections",
        body_lines=[
            "FY2025 target: $12M ARR — 3× YoY growth",
            "Gross margin: 78% at scale",
            "CAC payback: 14 months (improving to 10 months by Q4)",
            "Series B runway: 24 months post-raise",
            "Path to profitability: Q3 2026",
        ],
        title_color=BLACK,
    )

    # ---------- Slide 8: Call to Action ----------
    add_slide_with_title_content(
        prs,
        title_text="Next Steps",
        body_lines=[
            "Schedule a live demo: demo.novatech.io",
            "Review partnership prospectus (attached)",
            "Q&A and pilot agreement discussion",
            "Contact: partnerships@novatech.io | +1-415-555-0198",
        ],
        title_color=BLACK,
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
