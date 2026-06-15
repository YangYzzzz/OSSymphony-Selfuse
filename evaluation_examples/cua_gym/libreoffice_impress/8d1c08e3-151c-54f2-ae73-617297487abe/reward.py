"""
Reward Script: Duplicate slides 3 and 4, place duplicates at beginning of presentation
Task ID: impstruct_008
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): Total slide count is 7
  Component 2 (0.4): Slide title order matches expected sequence
  Component 3 (0.3): Duplicated slides content matches originals
"""

import os

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'impstruct_008'

# Expected slide title order after task completion
EXPECTED_TITLES = [
    'Solution',
    'Market Size',
    'Company Overview',
    'Problem',
    'Solution',
    'Market Size',
    'Ask',
]


def get_slide_title(slide):
    """Extract the title text from a slide."""
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        return slide.shapes.title.text.strip()
    # Fallback: look for any shape with text
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                return text
    return ''


def get_slide_texts(slide):
    """Get all non-empty text from a slide, for content comparison."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return texts


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)
    num_slides = len(slides)

    # Component 1: Total slide count is 7 (0.3 points)
    # Initial has 5 slides; after duplicating 2 slides we expect 7
    try:
        if num_slides == 7:
            print(f"PASS: Component 1 - Slide count is 7 (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 - Expected 7 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Slide title order matches expected sequence (0.4 points)
    # Expected: Solution, Market Size, Company Overview, Problem, Solution, Market Size, Ask
    try:
        actual_titles = [get_slide_title(s) for s in slides]
        if len(actual_titles) == len(EXPECTED_TITLES):
            matches = sum(1 for a, e in zip(actual_titles, EXPECTED_TITLES) if a == e)
            if matches == len(EXPECTED_TITLES):
                print(f"PASS: Component 2 - All 7 slide titles in correct order (0.4 pts)")
                total_score += 0.4
            else:
                print(f"FAIL: Component 2 - {matches}/{len(EXPECTED_TITLES)} titles match")
                print(f"  Expected: {EXPECTED_TITLES}")
                print(f"  Actual:   {actual_titles}")
        else:
            print(f"FAIL: Component 2 - Title count mismatch ({len(actual_titles)} vs {len(EXPECTED_TITLES)})")
            print(f"  Actual titles: {actual_titles}")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Duplicated slides have matching content with originals (0.3 points)
    # Slide 1 (copy of Solution) should match slide 5 (original Solution)
    # Slide 2 (copy of Market Size) should match slide 6 (original Market Size)
    try:
        if num_slides >= 7:
            copy1_texts = get_slide_texts(slides[0])
            orig1_texts = get_slide_texts(slides[4])
            copy2_texts = get_slide_texts(slides[1])
            orig2_texts = get_slide_texts(slides[5])

            pair1_match = copy1_texts == orig1_texts
            pair2_match = copy2_texts == orig2_texts

            if pair1_match and pair2_match:
                print(f"PASS: Component 3 - Duplicated slides match originals (0.3 pts)")
                total_score += 0.3
            else:
                if not pair1_match:
                    print(f"FAIL: Component 3 - Slide 1 (copy) text differs from slide 5 (original)")
                    print(f"  Copy texts:     {copy1_texts}")
                    print(f"  Original texts: {orig1_texts}")
                if not pair2_match:
                    print(f"FAIL: Component 3 - Slide 2 (copy) text differs from slide 6 (original)")
                    print(f"  Copy texts:     {copy2_texts}")
                    print(f"  Original texts: {orig2_texts}")
        else:
            print(f"FAIL: Component 3 - Not enough slides ({num_slides}) to compare duplicates")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/pitch_deck.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
