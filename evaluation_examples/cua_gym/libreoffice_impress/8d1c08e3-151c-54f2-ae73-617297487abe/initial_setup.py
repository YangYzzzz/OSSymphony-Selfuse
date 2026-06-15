"""
Initial Setup: Create a 5-slide pitch deck presentation
Task ID: impstruct_008
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
OUTPUT = f'{WORKDIR}/pitch_deck.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_content_slide(prs, title_text, body_lines):
    """Add a slide with title and bullet-point body content."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content layout
    slide.shapes.title.text = title_text

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Company Overview ---
    add_content_slide(prs, "Company Overview", [
        "NovaTech Solutions - Founded 2019 in San Francisco",
        "AI-powered supply chain optimization platform",
        "Team of 45 engineers, data scientists, and domain experts",
        "Serving 120+ enterprise clients across 8 countries",
        "Series B funded ($32M raised to date)",
    ])

    # --- Slide 2: Problem ---
    add_content_slide(prs, "Problem", [
        "Global supply chains lose $1.8T annually to inefficiency",
        "Legacy ERP systems lack real-time predictive capabilities",
        "Manual demand forecasting has 35-50% error rates",
        "Inventory carrying costs average 25% of product value",
        "Disruptions take 2-4 weeks to detect with current tools",
    ])

    # --- Slide 3: Solution ---
    add_content_slide(prs, "Solution", [
        "Proprietary ML engine processes 50M+ data points daily",
        "Real-time anomaly detection reduces disruption response to 4 hours",
        "Demand forecasting accuracy improved to 92% (vs. 55% industry avg)",
        "Automated reorder point optimization across multi-tier networks",
        "Seamless integration with SAP, Oracle, and Microsoft Dynamics",
    ])

    # --- Slide 4: Market Size ---
    add_content_slide(prs, "Market Size", [
        "Total Addressable Market (TAM): $48B by 2027",
        "Serviceable Available Market (SAM): $12B in target verticals",
        "Serviceable Obtainable Market (SOM): $1.2B in Year 5",
        "Growing at 14.3% CAGR driven by digital transformation",
        "Key verticals: Manufacturing, Retail, Pharma, Automotive",
    ])

    # --- Slide 5: Ask ---
    add_content_slide(prs, "Ask", [
        "Raising $18M Series C to accelerate growth",
        "Target: 300 enterprise clients by end of 2026",
        "Expand into APAC and Latin American markets",
        "Invest in GenAI copilot for procurement teams",
        "Projected ARR of $45M within 18 months post-funding",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Verify slide count
    verify = Presentation(OUTPUT)
    print(f'Slide count: {len(verify.slides)}')
    for i, slide in enumerate(verify.slides):
        title = slide.shapes.title.text if slide.shapes.title else "(no title)"
        print(f'  Slide {i+1}: {title}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
