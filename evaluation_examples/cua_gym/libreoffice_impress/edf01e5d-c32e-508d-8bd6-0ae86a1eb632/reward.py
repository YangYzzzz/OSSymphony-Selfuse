"""
FINAL REWARD SCRIPT - SUCCESS
Task: Convert comma-separated values in paragraph 3 into a table (comma delimiter).
Generated: 2025-10-17 12:44:32
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
import re
from pptx import Presentation


def verify_task(file_path: str) -> float:
    """Verify that paragraph 3's comma-separated values have been converted into
    a table (comma delimiter) inside the provided presentation file.

    Progressive scoring (max 1.0):
    1. A table exists in the presentation ................................ 0.5
    2. A row in that table contains the exact expected tokens ............ 0.3
    3. The original comma-separated paragraph is no longer present ........ 0.2
    """

    print(f"Verifying presentation: {file_path}")

    # Expected tokens that should now appear as individual cells
    expected_tokens = ["Apple", "Banana", "Cherry", "Date"]

    # ---------- Prerequisite: file exists & loads (NO points for this) ----------
    if not os.path.exists(file_path):
        print("✗ File not found – verification failed")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as exc:
        print(f"✗ Could not open presentation: {exc}")
        return 0.0

    total_score = 0.0  # progressive scoring accumulator

    # ---------- Requirement 1: A table exists -----------------------------------
    tables = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                tables.append(shape.table)

    if tables:
        print(f"✓ Found {len(tables)} table(s) in presentation (0.5 points)")
        total_score += 0.5
    else:
        print("✗ No table found – cannot satisfy task requirements")
        print(f"REWARD: {total_score}")
        return total_score  # cannot continue meaningful verification

    # ---------- Requirement 2: Correct row tokens ------------------------------
    row_with_tokens_found = False
    for tbl in tables:
        for row in tbl.rows:
            # Gather trimmed text for each cell in this row
            cell_texts = [cell.text_frame.text.strip() for cell in row.cells]

            # Compare only the length of expected tokens
            candidate = cell_texts[: len(expected_tokens)]
            if [t.lower() for t in candidate] == [et.lower() for et in expected_tokens]:
                row_with_tokens_found = True
                break
        if row_with_tokens_found:
            break

    if row_with_tokens_found:
        print("✓ Table row matches expected tokens (0.3 points)")
        total_score += 0.3
    else:
        print("✗ No table row matches the expected tokens")

    # ---------- Requirement 3: Original comma-separated paragraph removed ------
    # Aggregate all text from text-containing shapes for a simple presence check
    all_slide_text = "\n".join(
        shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
    )

    original_pattern = re.compile(r"Apple\s*,\s*Banana\s*,\s*Cherry\s*,\s*Date", re.IGNORECASE)
    if original_pattern.search(all_slide_text):
        print("✗ Original comma-separated list still present – no points for removal")
    else:
        print("✓ Original comma-separated list removed (0.2 points)")
        total_score += 0.2

    # Cap score at 1.0 and output results
    final_score = min(total_score, 1.0)
    print(f"Total Score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score


if __name__ == "__main__":
    # Path to the presentation the grader should evaluate
    PRESENTATION_PATH = "/home/user/convert_comma_separated_values_in_paragraph_3_into_a_table_comma_delimiter.pptx"
    verify_task(PRESENTATION_PATH)

