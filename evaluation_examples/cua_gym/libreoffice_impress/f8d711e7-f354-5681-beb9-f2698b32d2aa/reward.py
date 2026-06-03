"""
FINAL REWARD SCRIPT - SUCCESS
Task: I'm working on a presentation and want to spice up slide 4 a bit. Can someone guide me on how to draw a right-angled triangle and color it orange? I'm sure there's a way in LibreOffice Impress, just need a little help. Thanks!
Generated: 2025-08-07 12:12:42
Status: success
Model: o4-mini
Total Steps: 1
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor


def verify_task(file_path):
    print("Checking task completion for drawing right-angled orange triangle on slide 4")
    score = 0.0
    max_score = 1.0

    # 1. Verify the presentation file exists
    if os.path.exists(file_path):
        print(f"✓ File exists: {file_path} (0.1 points)")
        score += 0.1
    else:
        print(f"✗ File not found: {file_path}")
        print(f"REWARD: {score}")
        return score

    # 2. Load the presentation
    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded with {len(prs.slides)} slides (0.1 points)")
        score += 0.1
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        print(f"REWARD: {score}")
        return score

    # 3. Verify at least 4 slides
    slide_count = len(prs.slides)
    if slide_count >= 4:
        print(f"✓ Slide count >= 4 (0.1 points), found {slide_count}")
        score += 0.1
    else:
        print(f"✗ Slide count < 4, found {slide_count}")
        print(f"REWARD: {score}")
        return score

    # 4. Access slide 4 (index 3)
    try:
        slide4 = prs.slides[3]
    except Exception as e:
        print(f"✗ Cannot access slide 4: {e}")
        print(f"REWARD: {score}")
        return score

    # 5. Search for a right-angled triangle shape
    triangle_found = False
    for shape in slide4.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and \
           getattr(shape, 'auto_shape_type', None) == MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE:
            triangle_found = True
            print("✓ Found right-angled triangle shape (0.3 points)")
            score += 0.3

            # 6. Check that the triangle is filled with orange (RGB 255,165,0)
            fill = shape.fill
            rgb = None
            try:
                rgb = fill.fore_color.rgb
            except Exception:
                rgb = None

            expected_color = RGBColor(255, 165, 0)
            if rgb == expected_color:
                print("✓ Triangle colored orange correctly (0.5 points)")
                score += 0.5
            else:
                print(f"✗ Triangle color is {rgb}, expected {expected_color}")
            break

    if not triangle_found:
        print("✗ No right-angled triangle found on slide 4 (0 points for shape/color)")

    # 7. Final scoring
    final_score = min(score, max_score)
    print(f"Total score: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score


if __name__ == '__main__':
    pptx_path = '/home/user/im_working_on_a_presentation_and_want_to_spice_up_slide_4_a_bit_can_someone_guide_me_on_how_to_draw_.pptx'
    verify_task(pptx_path)

