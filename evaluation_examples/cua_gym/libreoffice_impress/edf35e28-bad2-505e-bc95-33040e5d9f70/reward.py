"""
FINAL REWARD SCRIPT - SUCCESS
Task: Set the line spacing of paragraphs 2–3 to Exactly 18 pt.
Generated: 2025-10-17 08:50:25
Status: success
Model: azure-o3
Total Steps: 6
"""

import os
import math
from pptx import Presentation

# ----------------------------------------------------------------------------
# Reward Script : Verify that paragraphs 2 and 3 in the presentation have
# their line-spacing set to *Exactly* 18 pt.
# ----------------------------------------------------------------------------
# Scoring logic (progressive):
#   • 0.5 points if paragraph 2 is formatted correctly
#   • 0.5 points if paragraph 3 is formatted correctly
#   • Total score is capped at 1.0
# ----------------------------------------------------------------------------
# IMPORTANT:  No points are awarded for file existence or successful loading –
# those are prerequisites, not accomplishments.
# ----------------------------------------------------------------------------

EMUS_PER_PT = 12700  # 1 point = 12700 English Metric Units (EMU)
EXPECTED_PT  = 18.0
EXPECTED_EMU = EXPECTED_PT * EMUS_PER_PT
TOLERANCE    = 1000  # ≈ 0.08 pt – small safety margin

PARAGRAPHS_TO_CHECK = {2, 3}  # 1-based indices


def verify_line_spacing(file_path: str) -> float:
    """Verify that paragraphs 2 and 3 have line-spacing == 18 pt.

    Returns
    -------
    float
        Reward between 0.0 and 1.0
    """

    print(f"Verifying presentation: {file_path}\n")

    # ---------- Fail-fast checks (no score awarded) ----------
    if not os.path.exists(file_path):
        print("✗ File not found – task not completed")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Unable to open PPTX: {e}")
        print("REWARD: 0.0")
        return 0.0

    # ---------- Actual verification ----------
    correctly_formatted: set[int] = set()  # which paragraph indices passed

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if not shape.has_text_frame:
                continue

            tf = shape.text_frame
            for p_idx, para in enumerate(tf.paragraphs, start=1):
                # Only evaluate the target paragraph numbers (2 and 3)
                if p_idx not in PARAGRAPHS_TO_CHECK:
                    continue

                raw_ls = para.line_spacing  # typically an int (EMU) or None

                # Display helper info for transparency
                if raw_ls is None:
                    print(f"Slide {slide_idx}, Shape {shape_idx}, Paragraph {p_idx}: "
                          f"line_spacing=None (expected {EXPECTED_PT} pt)")
                    continue

                # Ensure we are working with a numeric EMU value
                if hasattr(raw_ls, "pt"):
                    ls_emu = raw_ls.pt * EMUS_PER_PT
                else:
                    ls_emu = raw_ls

                ls_pt = ls_emu / EMUS_PER_PT
                print(f"Slide {slide_idx}, Shape {shape_idx}, Paragraph {p_idx}: "
                      f"line_spacing={ls_pt:.2f} pt (raw {ls_emu})")

                # Compare with tolerance
                if math.isclose(ls_emu, EXPECTED_EMU, abs_tol=TOLERANCE):
                    print("  ✓ Line spacing matches expected value")
                    correctly_formatted.add(p_idx)
                else:
                    print("  ✗ Line spacing does NOT match expected value")

    # ---------- Scoring ----------
    score = 0.0
    if 2 in correctly_formatted:
        score += 0.5
    if 3 in correctly_formatted:
        score += 0.5

    score = min(score, 1.0)  # safety cap

    print(f"\nTotal score: {score}/1.0")
    print(f"REWARD: {score}")
    return score


# -------------------------------------------------------------------------
# Execute verification when the script is run
# -------------------------------------------------------------------------
if __name__ == "__main__":
    PRESENTATION_PATH = "/home/user/set_the_line_spacing_of_paragraphs_23_to_exactly_18_pt.pptx"
    verify_line_spacing(PRESENTATION_PATH)

