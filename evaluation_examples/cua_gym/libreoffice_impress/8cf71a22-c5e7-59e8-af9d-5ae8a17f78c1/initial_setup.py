"""
Initial Setup: Build a 12-slide sales pitch presentation with no navigation buttons.
Task ID: impress_sales_060
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_060'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a styled text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    # Standard widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ========== Slide 1: Title Slide ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    bg = slide1.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x44)

    add_textbox(slide1, Inches(2), Inches(1.5), Inches(9), Inches(1.5),
                "TechVista Solutions", font_size=40, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(2), Inches(3.0), Inches(9), Inches(1.0),
                "Enterprise Cloud Platform — Interactive Sales Pitch",
                font_size=22, color=RGBColor(0xBB, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(3), Inches(4.5), Inches(7), Inches(0.8),
                "Q2 2025 | Prepared for Strategic Partners",
                font_size=16, color=RGBColor(0x99, 0xAA, 0xBB), alignment=PP_ALIGN.CENTER)

    # ========== Slide 2: Company Overview ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide2, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                "Company Overview", font_size=32, bold=True,
                color=RGBColor(0x2B, 0x6C, 0xB0))
    add_textbox(slide2, Inches(0.8), Inches(1.6), Inches(11), Inches(4.5),
                "TechVista Solutions was founded in 2018 by former AWS and Azure engineers "
                "who saw a gap in the market for mid-enterprise cloud migration. Headquartered "
                "in Austin, TX with offices in London and Singapore, we serve over 340 clients "
                "across 28 countries.\n\n"
                "Our mission: Make enterprise cloud adoption seamless, secure, and cost-effective.\n\n"
                "Key milestones:\n"
                "• 2019: First 50 enterprise clients onboarded\n"
                "• 2021: Series B funding — $45M led by Sequoia Capital\n"
                "• 2023: Named a Gartner Cool Vendor in Cloud Management\n"
                "• 2024: Revenue reached $128M ARR with 97% retention rate",
                font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # ========== Slide 3: Team & Leadership ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide3, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                "Leadership Team", font_size=32, bold=True,
                color=RGBColor(0x2B, 0x6C, 0xB0))
    add_textbox(slide3, Inches(0.8), Inches(1.6), Inches(11), Inches(4.5),
                "CEO: Dr. Ananya Patel — Former VP Engineering at AWS, MIT PhD\n"
                "CTO: Marcus Chen — Ex-Google Cloud Architect, 15+ years distributed systems\n"
                "CFO: Sarah Williams — Previously at Deloitte, IPO experience with 3 tech firms\n"
                "VP Sales: James Rodriguez — Built $200M pipeline at Salesforce\n"
                "VP Engineering: Mei Lin — Former Microsoft Azure lead, 40+ patents",
                font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # ========== Slide 4: Platform Features Overview ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide4, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                "Platform Features", font_size=32, bold=True,
                color=RGBColor(0x2B, 0x6C, 0xB0))
    add_textbox(slide4, Inches(0.8), Inches(1.6), Inches(11), Inches(4.5),
                "CloudBridge™ unifies migration, monitoring, and optimization in one platform:\n\n"
                "• Automated Migration Engine — Zero-downtime workload transfers\n"
                "• Real-Time Cost Optimizer — AI-driven spend reduction averaging 34%\n"
                "• Security Command Center — SOC2, HIPAA, GDPR compliance built-in\n"
                "• Multi-Cloud Orchestrator — AWS, Azure, GCP from a single pane\n"
                "• Custom Integration Hub — 200+ pre-built connectors for enterprise tools",
                font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # ========== Slide 5: Technical Architecture ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide5, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                "Technical Architecture", font_size=32, bold=True,
                color=RGBColor(0x2B, 0x6C, 0xB0))
    add_textbox(slide5, Inches(0.8), Inches(1.6), Inches(11), Inches(4.5),
                "Three-tier architecture designed for enterprise scale:\n\n"
                "Presentation Layer: React-based dashboard, mobile-responsive, SSO integration\n"
                "Service Layer: Kubernetes-orchestrated microservices, auto-scaling to 10K+ nodes\n"
                "Data Layer: Distributed storage with 99.999% uptime SLA, encrypted at rest & transit\n\n"
                "Performance benchmarks:\n"
                "• Migration throughput: 50TB/hour sustained\n"
                "• API response time: < 50ms p99\n"
                "• Platform uptime: 99.98% over trailing 12 months",
                font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # ========== Slide 6: Security & Compliance ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide6, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                "Security & Compliance", font_size=32, bold=True,
                color=RGBColor(0x2B, 0x6C, 0xB0))
    add_textbox(slide6, Inches(0.8), Inches(1.6), Inches(11), Inches(4.5),
                "Enterprise-grade security at every layer:\n\n"
                "• SOC 2 Type II certified annually since 2021\n"
                "• HIPAA BAA available for healthcare clients\n"
                "• GDPR compliant with EU data residency options\n"
                "• FedRAMP Moderate authorization (in progress)\n"
                "• Zero-trust network architecture with mTLS\n"
                "• 24/7 SOC with < 15 min incident response SLA\n\n"
                "Annual penetration testing by CrowdStrike; zero critical findings in last audit.",
                font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # ========== Slide 7: Pricing Overview ==========
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide7, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                "Pricing Plans", font_size=32, bold=True,
                color=RGBColor(0x2B, 0x6C, 0xB0))
    add_textbox(slide7, Inches(0.8), Inches(1.6), Inches(11), Inches(4.5),
                "Flexible pricing to match your growth stage:\n\n"
                "Starter — $2,500/month\n"
                "  Up to 50 workloads, basic monitoring, email support\n\n"
                "Professional — $8,500/month\n"
                "  Up to 200 workloads, full optimization suite, dedicated CSM\n\n"
                "Enterprise — Custom pricing\n"
                "  Unlimited workloads, premium SLA, on-site deployment option\n\n"
                "All plans include: 30-day free trial, no long-term contracts, 99.9% SLA",
                font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # ========== Slide 8: ROI Analysis ==========
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide8, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                "ROI Analysis", font_size=32, bold=True,
                color=RGBColor(0x2B, 0x6C, 0xB0))
    add_textbox(slide8, Inches(0.8), Inches(1.6), Inches(11), Inches(4.5),
                "Average client results within 12 months:\n\n"
                "• Cloud infrastructure costs: Reduced 34% ($420K avg savings)\n"
                "• Migration timeline: 60% faster than manual approaches\n"
                "• IT staff productivity: 25% increase through automation\n"
                "• Downtime incidents: 78% reduction\n"
                "• Compliance audit preparation: From 6 weeks to 3 days\n\n"
                "Typical payback period: 4.2 months\n"
                "3-year TCO advantage: 2.8x vs. in-house tooling",
                font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # ========== Slide 9: Case Study — HealthFirst ==========
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide9, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                "Case Study: HealthFirst Medical Group", font_size=32, bold=True,
                color=RGBColor(0x2B, 0x6C, 0xB0))
    add_textbox(slide9, Inches(0.8), Inches(1.6), Inches(11), Inches(4.5),
                "Challenge: Migrate 2,400 workloads from on-prem to AWS while maintaining "
                "HIPAA compliance for 1.2M patient records.\n\n"
                "Solution: CloudBridge automated migration with zero-downtime cutover, "
                "built-in HIPAA controls, and real-time compliance dashboard.\n\n"
                "Results:\n"
                "• Migration completed in 8 weeks (estimated 6 months manually)\n"
                "• $1.8M annual infrastructure savings\n"
                "• Zero data breaches, passed all compliance audits\n"
                "• 99.99% application uptime during migration",
                font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # ========== Slide 10: Case Study — RetailMax ==========
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide10, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                "Case Study: RetailMax International", font_size=32, bold=True,
                color=RGBColor(0x2B, 0x6C, 0xB0))
    add_textbox(slide10, Inches(0.8), Inches(1.6), Inches(11), Inches(4.5),
                "Challenge: Multi-cloud strategy across AWS and Azure for 15 countries, "
                "with strict data sovereignty requirements in EU and APAC.\n\n"
                "Solution: CloudBridge Multi-Cloud Orchestrator with regional data residency "
                "policies and unified cost management across both providers.\n\n"
                "Results:\n"
                "• 42% reduction in total cloud spend ($3.2M annually)\n"
                "• Single dashboard for 3,800+ workloads across 2 clouds\n"
                "• Automated compliance reporting for 7 regulatory frameworks\n"
                "• Deployment time reduced from 4 hours to 12 minutes",
                font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # ========== Slide 11: Client Testimonials ==========
    slide11 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide11, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                "What Our Clients Say", font_size=32, bold=True,
                color=RGBColor(0x2B, 0x6C, 0xB0))
    add_textbox(slide11, Inches(0.8), Inches(1.6), Inches(11), Inches(4.5),
                '"TechVista transformed our cloud strategy. What we thought would take '
                'a year was done in weeks." — Dr. Rebecca Torres, CIO, HealthFirst\n\n'
                '"The cost savings alone justified the investment within the first quarter. '
                'The platform practically pays for itself." — Priya Sharma, VP Ops, RetailMax\n\n'
                '"Finally, a tool that lets us manage AWS and Azure without separate teams. '
                'Game changer." — David Park, Cloud Architect, NovaTech Industries',
                font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # ========== Slide 12: Contact Us ==========
    slide12 = prs.slides.add_slide(prs.slide_layouts[5])
    bg12 = slide12.background.fill
    bg12.solid()
    bg12.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x44)
    add_textbox(slide12, Inches(2), Inches(1.0), Inches(9), Inches(1.2),
                "Let's Get Started", font_size=36, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_textbox(slide12, Inches(2), Inches(2.8), Inches(9), Inches(3.5),
                "Contact our sales team to schedule a personalized demo:\n\n"
                "Email: sales@techvista.io\n"
                "Phone: +1 (512) 555-0192\n"
                "Web: www.techvista.io/demo\n\n"
                "James Rodriguez, VP Sales\n"
                "james.rodriguez@techvista.io | Direct: +1 (512) 555-0234",
                font_size=18, color=RGBColor(0xCC, 0xDD, 0xEE), alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
