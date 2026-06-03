"""
Reward Script: Interactive navigation system with action buttons
Task ID: impress_sales_060
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Slide 1 has 5 navigation buttons with correct labels
  Component 2 (0.25): Navigation buttons link to correct target slides
  Component 3 (0.20): Navigation buttons styled correctly (fill #2B6CB0, white text, ~12pt)
  Component 4 (0.25): Slides 2-12 each have a 'Home' button linking to slide 1
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_060'

# Expected nav buttons on slide 1: label -> target slide number (1-indexed)
EXPECTED_NAV_BUTTONS = {
    'Overview': 2,
    'Features': 4,
    'Pricing': 7,
    'Case Studies': 9,
    'Contact': 12,
}


def get_slide_rels(pptx_path, slide_num):
    """Get relationship ID -> target slide number mapping for a given slide."""
    rels = {}
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        rels_path = f'ppt/slides/_rels/slide{slide_num}.xml.rels'
        try:
            with zf.open(rels_path) as f:
                root = ET.parse(f).getroot()
                for rel in root:
                    target = rel.attrib.get('Target', '')
                    rel_type = rel.attrib.get('Type', '')
                    rel_id = rel.attrib.get('Id', '')
                    if 'relationships/slide' in rel_type and target.startswith('slide'):
                        # Extract slide number from e.g. 'slide4.xml'
                        try:
                            target_num = int(target.replace('slide', '').replace('.xml', ''))
                            rels[rel_id] = target_num
                        except ValueError:
                            pass
        except KeyError:
            pass
    return rels


def get_shape_hyperlink_rel_id(shape):
    """Extract the relationship ID from a shape's hyperlink click action."""
    from pptx.oxml.ns import qn
    el = shape._element
    hlink = el.find('.//' + qn('a:hlinkClick'))
    if hlink is not None:
        r_id = hlink.attrib.get(qn('r:id'), '')
        action = hlink.attrib.get('action', '')
        if 'hlinksldjump' in action:
            return r_id
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Pt
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have 12 slides
    if len(prs.slides) != 12:
        print(f"PRECONDITION FAIL: Expected 12 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide1 = prs.slides[0]

    # =========================================================================
    # Component 1: Slide 1 has 5 navigation buttons with correct labels (0.30)
    # =========================================================================
    try:
        # Find all auto_shape (rounded rectangle) shapes on slide 1 that have text
        nav_buttons = {}
        for shape in slide1.shapes:
            if hasattr(shape, 'text_frame') and shape.shape_type is not None:
                # Check if it's an auto shape (type 1) with a hyperlink
                shape_text = shape.text.strip()
                hlink_rid = get_shape_hyperlink_rel_id(shape)
                if hlink_rid and shape_text:
                    nav_buttons[shape_text] = {'shape': shape, 'rel_id': hlink_rid}

        # Check all 5 expected labels are present
        found_labels = set(nav_buttons.keys())
        expected_labels = set(EXPECTED_NAV_BUTTONS.keys())
        matching_labels = found_labels & expected_labels
        label_count = len(matching_labels)

        if label_count == 5:
            print(f"PASS: Component 1 — All 5 nav buttons found on slide 1: {sorted(matching_labels)} (0.30 pts)")
            total_score += 0.30
        elif label_count >= 3:
            partial = round(0.30 * label_count / 5, 2)
            print(f"PARTIAL: Component 1 — {label_count}/5 nav buttons found: {sorted(matching_labels)} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — Only {label_count}/5 expected nav buttons found. Found: {sorted(found_labels)}, expected: {sorted(expected_labels)}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # =========================================================================
    # Component 2: Navigation buttons link to correct target slides (0.25)
    # =========================================================================
    try:
        slide1_rels = get_slide_rels(file_path, 1)
        correct_links = 0
        for label, target_slide in EXPECTED_NAV_BUTTONS.items():
            if label in nav_buttons:
                rel_id = nav_buttons[label]['rel_id']
                actual_target = slide1_rels.get(rel_id)
                if actual_target == target_slide:
                    correct_links += 1
                    print(f"  Link OK: '{label}' -> slide {actual_target}")
                else:
                    print(f"  Link WRONG: '{label}' -> slide {actual_target} (expected {target_slide})")
            else:
                print(f"  Link MISSING: '{label}' button not found")

        if correct_links == 5:
            print(f"PASS: Component 2 — All 5 nav buttons link to correct slides (0.25 pts)")
            total_score += 0.25
        elif correct_links >= 1:
            partial = round(0.25 * correct_links / 5, 2)
            print(f"PARTIAL: Component 2 — {correct_links}/5 correct links ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No correct hyperlink targets")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # =========================================================================
    # Component 3: Button styling — fill #2B6CB0, white text, ~12pt (0.20)
    # =========================================================================
    try:
        styled_count = 0
        for label in EXPECTED_NAV_BUTTONS:
            if label not in nav_buttons:
                continue
            shape = nav_buttons[label]['shape']
            checks_passed = 0
            total_checks = 3  # fill color, font color, font size

            # Check fill color
            try:
                if shape.fill.type is not None and str(shape.fill.fore_color.rgb).upper() == '2B6CB0':
                    checks_passed += 1
                else:
                    print(f"  Style '{label}': fill color mismatch (found {shape.fill.fore_color.rgb})")
            except Exception:
                print(f"  Style '{label}': could not read fill color")

            # Check font color and size from first run
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    # White text
                    try:
                        if run.font.color.type is not None and str(run.font.color.rgb).upper() == 'FFFFFF':
                            checks_passed += 1
                        else:
                            print(f"  Style '{label}': font color mismatch (found {run.font.color.rgb})")
                    except Exception:
                        print(f"  Style '{label}': could not read font color")

                    # ~12pt (152400 EMU = 12pt; allow some tolerance)
                    if run.font.size is not None:
                        size_pt = run.font.size / 12700  # EMU to pt
                        if 11.0 <= size_pt <= 13.0:
                            checks_passed += 1
                        else:
                            print(f"  Style '{label}': font size {size_pt}pt (expected ~12pt)")
                    break  # only check first run
                break  # only check first paragraph

            if checks_passed == total_checks:
                styled_count += 1

        if styled_count == 5:
            print(f"PASS: Component 3 — All 5 buttons styled correctly: fill #2B6CB0, white 12pt (0.20 pts)")
            total_score += 0.20
        elif styled_count >= 1:
            partial = round(0.20 * styled_count / 5, 2)
            print(f"PARTIAL: Component 3 — {styled_count}/5 buttons styled correctly ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No buttons have correct styling")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # =========================================================================
    # Component 4: Slides 2-12 each have a 'Home' button linking to slide 1 (0.25)
    # =========================================================================
    try:
        home_count = 0
        total_expected = 11  # slides 2 through 12

        for slide_idx in range(1, 12):  # 0-indexed: slides 2-12
            slide = prs.slides[slide_idx]
            slide_num = slide_idx + 1
            slide_rels = get_slide_rels(file_path, slide_num)

            found_home = False
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    shape_text = shape.text.strip().lower()
                    if shape_text == 'home':
                        hlink_rid = get_shape_hyperlink_rel_id(shape)
                        if hlink_rid:
                            target = slide_rels.get(hlink_rid)
                            if target == 1:
                                found_home = True
                                break

            if found_home:
                home_count += 1
            else:
                print(f"  Home MISSING on slide {slide_num}")

        if home_count == total_expected:
            print(f"PASS: Component 4 — All 11 slides (2-12) have Home button linking to slide 1 (0.25 pts)")
            total_score += 0.25
        elif home_count >= 1:
            partial = round(0.25 * home_count / total_expected, 2)
            print(f"PARTIAL: Component 4 — {home_count}/{total_expected} slides have correct Home button ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — No slides have a Home button linking to slide 1")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Main entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'

persist_app_state()

if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
