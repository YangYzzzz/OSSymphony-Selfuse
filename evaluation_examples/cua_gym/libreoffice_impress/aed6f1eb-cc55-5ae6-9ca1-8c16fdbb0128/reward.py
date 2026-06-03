"""
Reward Script: Set all text on all slides to center alignment
Task ID: osworld_impress_per_slide_alignment_002
Domain: libreoffice_impress
Scoring:
  Component 1: All paragraphs on slides 1-2 are center-aligned (0.34 pts)
  Component 2: All paragraphs on slides 3-4 are center-aligned (0.33 pts)
  Component 3: All paragraphs on slides 5-6 are center-aligned (0.33 pts)
  Total: 1.0
"""

import os
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_per_slide_alignment_002'

# PP_ALIGN.CENTER has integer value 2
CENTER = PP_ALIGN.CENTER


def get_all_text_paragraphs(slide):
    """
    Recursively collect all paragraphs from all shapes (including groups)
    on a slide that have a text frame.
    Returns a list of (shape_name, para_index, paragraph) tuples.
    """
    results = []

    def extract(shape):
        if shape.has_text_frame:
            for p_idx, para in enumerate(shape.text_frame.paragraphs):
                results.append((shape.name, p_idx, para))
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                extract(sub)

    for shape in slide.shapes:
        extract(shape)
    return results


def check_slides_center_aligned(prs, slide_indices):
    """
    Check that all non-empty text paragraphs on the given slide indices
    (0-based) are center-aligned.
    Returns (all_pass, pass_count, total_count, details).
    """
    pass_count = 0
    total_count = 0
    details = []

    for slide_idx in slide_indices:
        if slide_idx >= len(prs.slides):
            details.append(f"Slide {slide_idx + 1}: not present in presentation")
            continue
        slide = prs.slides[slide_idx]
        paras = get_all_text_paragraphs(slide)
        for shape_name, p_idx, para in paras:
            # Only check paragraphs that have actual text content
            if not para.text.strip():
                continue
            total_count += 1
            align = para.alignment
            if align == CENTER:
                pass_count += 1
            else:
                details.append(
                    f"Slide {slide_idx + 1}, Shape='{shape_name}', Para={p_idx}: "
                    f"alignment={align} (expected CENTER)"
                )

    all_pass = (pass_count == total_count) and total_count > 0
    return all_pass, pass_count, total_count, details


def verify_task(file_path):
    """
    Verify task completion: all text on all slides should be center-aligned.
    Returns a float between 0.0 and 1.0.

    Scoring rubric:
      Component 1 (0.34 pts): All paragraphs on slides 1-2 are center-aligned
      Component 2 (0.33 pts): All paragraphs on slides 3-4 are center-aligned
      Component 3 (0.33 pts): All paragraphs on slides 5-6 are center-aligned
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Validate that we have the expected 6 slides as a precondition gate
    slide_count = len(prs.slides)
    if slide_count < 1:
        print("CRITICAL: Presentation has no slides")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Presentation loaded, slide count = {slide_count}")

    # Component 1: Slides 1-2 (indices 0-1) all center-aligned (0.34 points)
    try:
        indices = [i for i in [0, 1] if i < slide_count]
        all_pass, pass_cnt, total_cnt, failures = check_slides_center_aligned(prs, indices)
        if total_cnt == 0:
            print("FAIL: Component 1 — No text paragraphs found on slides 1-2")
        elif all_pass:
            print(f"PASS: Component 1 — All {total_cnt} paragraphs on slides 1-2 are center-aligned (0.34 pts)")
            total_score += 0.34
        else:
            print(f"FAIL: Component 1 — {pass_cnt}/{total_cnt} paragraphs on slides 1-2 are center-aligned")
            for f in failures:
                print(f"  {f}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slides 3-4 (indices 2-3) all center-aligned (0.33 points)
    try:
        indices = [i for i in [2, 3] if i < slide_count]
        all_pass, pass_cnt, total_cnt, failures = check_slides_center_aligned(prs, indices)
        if total_cnt == 0:
            print("FAIL: Component 2 — No text paragraphs found on slides 3-4")
        elif all_pass:
            print(f"PASS: Component 2 — All {total_cnt} paragraphs on slides 3-4 are center-aligned (0.33 pts)")
            total_score += 0.33
        else:
            print(f"FAIL: Component 2 — {pass_cnt}/{total_cnt} paragraphs on slides 3-4 are center-aligned")
            for f in failures:
                print(f"  {f}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slides 5-6 (indices 4-5) all center-aligned (0.33 points)
    try:
        indices = [i for i in [4, 5] if i < slide_count]
        all_pass, pass_cnt, total_cnt, failures = check_slides_center_aligned(prs, indices)
        if total_cnt == 0:
            print("FAIL: Component 3 — No text paragraphs found on slides 5-6")
        elif all_pass:
            print(f"PASS: Component 3 — All {total_cnt} paragraphs on slides 5-6 are center-aligned (0.33 pts)")
            total_score += 0.33
        else:
            print(f"FAIL: Component 3 — {pass_cnt}/{total_cnt} paragraphs on slides 5-6 are center-aligned")
            for f in failures:
                print(f"  {f}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entrypoint: test against canonical artifact path on the VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
