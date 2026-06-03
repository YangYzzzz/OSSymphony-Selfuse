"""
Initial Setup: 6-slide product overview deck with all text left-aligned
Task ID: osworld_impress_per_slide_alignment_002
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_per_slide_alignment_002'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_paragraph_left(para, text):
    """Set a paragraph's text and explicitly set alignment to LEFT."""
    para.text = text
    para.alignment = PP_ALIGN.LEFT


def create_initial():
    prs = Presentation()
    # Use standard 16:9 slide dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
    title1 = slide1.shapes.title
    title1.text = "NovaTech Pro Series"
    for para in title1.text_frame.paragraphs:
        para.alignment = PP_ALIGN.LEFT

    subtitle1 = slide1.placeholders[1]
    subtitle1.text = "Next-Generation Enterprise Solutions\nQ2 2025 Product Overview"
    for para in subtitle1.text_frame.paragraphs:
        para.alignment = PP_ALIGN.LEFT

    # --- Slide 2: Company Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title2 = slide2.shapes.title
    title2.text = "Company Overview"
    for para in title2.text_frame.paragraphs:
        para.alignment = PP_ALIGN.LEFT

    body2 = slide2.placeholders[1]
    tf2 = body2.text_frame
    tf2.text = "Founded in 2011, NovaTech delivers enterprise-grade software solutions"
    for para in tf2.paragraphs:
        para.alignment = PP_ALIGN.LEFT
    p2a = tf2.add_paragraph()
    p2a.text = "Over 2,400 clients across 38 countries"
    p2a.alignment = PP_ALIGN.LEFT
    p2b = tf2.add_paragraph()
    p2b.text = "Annual revenue exceeding $340M in fiscal year 2024"
    p2b.alignment = PP_ALIGN.LEFT
    p2c = tf2.add_paragraph()
    p2c.text = "Headquartered in Austin, TX with regional offices in London, Singapore, and Toronto"
    p2c.alignment = PP_ALIGN.LEFT

    # --- Slide 3: Product Features ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = "Key Product Features"
    for para in title3.text_frame.paragraphs:
        para.alignment = PP_ALIGN.LEFT

    body3 = slide3.placeholders[1]
    tf3 = body3.text_frame
    tf3.text = "Real-time analytics dashboard with customizable KPI widgets"
    for para in tf3.paragraphs:
        para.alignment = PP_ALIGN.LEFT
    features = [
        "AI-powered predictive maintenance — reduces downtime by up to 47%",
        "End-to-end encryption with SOC 2 Type II and ISO 27001 compliance",
        "Seamless API integration with Salesforce, SAP, and Oracle ERP systems",
        "Automated reporting with scheduled PDF and Excel exports",
    ]
    for feat in features:
        p = tf3.add_paragraph()
        p.text = feat
        p.alignment = PP_ALIGN.LEFT

    # --- Slide 4: Pricing & Plans ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    title4.text = "Pricing & Plans"
    for para in title4.text_frame.paragraphs:
        para.alignment = PP_ALIGN.LEFT

    body4 = slide4.placeholders[1]
    tf4 = body4.text_frame
    tf4.text = "Starter: $299/month — Up to 10 users, core analytics, email support"
    for para in tf4.paragraphs:
        para.alignment = PP_ALIGN.LEFT
    plans = [
        "Professional: $799/month — Up to 50 users, advanced AI features, priority support",
        "Enterprise: Custom pricing — Unlimited users, dedicated account manager, SLA guarantee",
        "All plans include a 30-day free trial with no credit card required",
        "Volume discounts available for multi-year contracts",
    ]
    for plan in plans:
        p = tf4.add_paragraph()
        p.text = plan
        p.alignment = PP_ALIGN.LEFT

    # --- Slide 5: Customer Success Stories ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    title5.text = "Customer Success Stories"
    for para in title5.text_frame.paragraphs:
        para.alignment = PP_ALIGN.LEFT

    body5 = slide5.placeholders[1]
    tf5 = body5.text_frame
    tf5.text = "GlobalMed Healthcare — Reduced reporting time by 68% using automated dashboards"
    for para in tf5.paragraphs:
        para.alignment = PP_ALIGN.LEFT
    stories = [
        "Meridian Logistics — Achieved $2.1M annual savings through predictive fleet maintenance",
        "FinCore Bank — Passed regulatory audit with zero findings after deploying compliance module",
        "Apex Manufacturing — Increased OEE from 71% to 89% within six months of deployment",
    ]
    for story in stories:
        p = tf5.add_paragraph()
        p.text = story
        p.alignment = PP_ALIGN.LEFT

    # --- Slide 6: Next Steps & Contact ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    title6.text = "Next Steps"
    for para in title6.text_frame.paragraphs:
        para.alignment = PP_ALIGN.LEFT

    body6 = slide6.placeholders[1]
    tf6 = body6.text_frame
    tf6.text = "Schedule a personalized demo with your regional sales engineer"
    for para in tf6.paragraphs:
        para.alignment = PP_ALIGN.LEFT
    steps = [
        "Contact us: sales@novatech.com | +1 (800) 555-0192",
        "Visit our portal: www.novatech.com/enterprise",
        "Request a proof-of-concept deployment within 72 hours",
        "Join our quarterly webinar series every third Tuesday",
    ]
    for step in steps:
        p = tf6.add_paragraph()
        p.text = step
        p.alignment = PP_ALIGN.LEFT

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open file in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
