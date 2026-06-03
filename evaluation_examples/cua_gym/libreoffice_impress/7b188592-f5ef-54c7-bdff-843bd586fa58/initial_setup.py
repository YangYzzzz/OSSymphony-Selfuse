"""
Initial Setup: Create a 7-slide Data Science Intro presentation with slide 5 empty below title.
Task ID: impress_stu_028
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_028'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, text in enumerate(bullets):
        if i == 0:
            tf.paragraphs[0].text = text
        else:
            p = tf.add_paragraph()
            p.text = text
            p.level = 0
    return slide


def add_title_only_slide(prs, title_text):
    # Use layout 5 (Blank) and add a title textbox manually
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    add_title_slide(prs, "Introduction to Data Science",
                    "Fundamentals, Tools, and Workflows\nPrepared by Dr. Elena Rodriguez")

    # Slide 2: What is Data Science?
    add_content_slide(prs, "What is Data Science?", [
        "An interdisciplinary field combining statistics, computer science, and domain expertise",
        "Extracts knowledge and insights from structured and unstructured data",
        "Involves data collection, cleaning, analysis, and visualization",
        "Drives decision-making across industries including healthcare, finance, and technology",
        "Rapidly growing field with increasing demand for skilled practitioners",
    ])

    # Slide 3: Key Skills
    add_content_slide(prs, "Key Skills for Data Scientists", [
        "Programming: Python, R, SQL for data manipulation and analysis",
        "Statistics & Probability: Hypothesis testing, regression, Bayesian methods",
        "Machine Learning: Supervised and unsupervised algorithms, model evaluation",
        "Data Visualization: Creating compelling charts and dashboards",
        "Communication: Translating technical findings into business insights",
        "Domain Knowledge: Understanding industry-specific challenges and context",
    ])

    # Slide 4: Tools & Technologies
    add_content_slide(prs, "Tools & Technologies", [
        "Python Libraries: pandas, NumPy, scikit-learn, TensorFlow, PyTorch",
        "Visualization: Matplotlib, Seaborn, Plotly, Tableau, Power BI",
        "Big Data: Apache Spark, Hadoop, Kafka for large-scale processing",
        "Cloud Platforms: AWS SageMaker, Google AI Platform, Azure ML",
        "Version Control: Git and GitHub for collaborative development",
        "Notebooks: Jupyter, Google Colab for interactive experimentation",
    ])

    # Slide 5: Data Science Pipeline (TITLE ONLY - empty below)
    add_title_only_slide(prs, "Data Science Pipeline")

    # Slide 6: Case Studies
    add_content_slide(prs, "Real-World Case Studies", [
        "Netflix: Recommendation engine drives 80% of content watched",
        "Spotify: Discover Weekly uses collaborative filtering and NLP on playlists",
        "Uber: Dynamic pricing model optimizes supply-demand matching in real time",
        "Healthcare: Predictive models for early disease detection reduce mortality rates",
        "Retail: Customer segmentation enables personalized marketing campaigns",
    ])

    # Slide 7: Conclusion & Next Steps
    add_content_slide(prs, "Conclusion & Next Steps", [
        "Data science is transforming how organizations make decisions",
        "A structured pipeline ensures reproducible and reliable results",
        "Continuous learning is essential as the field evolves rapidly",
        "Next session: Hands-on workshop with real datasets",
        "Resources: course materials available on the shared drive",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
