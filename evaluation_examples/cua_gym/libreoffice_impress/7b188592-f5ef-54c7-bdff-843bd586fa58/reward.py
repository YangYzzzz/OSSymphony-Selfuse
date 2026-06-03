"""
Reward Script: Create flowchart on slide 5 with four rectangles connected by arrows
Task ID: impress_stu_028
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): Four rectangles with correct labels on slide 5
  Component 2 (0.3): Three connector/arrow lines between rectangles
  Component 3 (0.2): Horizontal arrangement in correct left-to-right order
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_028'


def persist_app_state(domain: str):
    """Attempt to save any unsaved LibreOffice edits."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 5 slides
    if len(prs.slides) < 5:
        print(f"PRECONDITION FAIL: Need at least 5 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # Slide 5 (0-indexed)

    # Collect rectangles and connectors on slide 5
    rectangles = []
    connectors = []
    for shape in slide.shapes:
        # Check for rectangle AUTO_SHAPE
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                from pptx.enum.shapes import MSO_SHAPE
                if shape.auto_shape_type == MSO_SHAPE.RECTANGLE:
                    text = ""
                    if shape.has_text_frame:
                        text = " ".join(
                            p.text.strip() for p in shape.text_frame.paragraphs
                        ).strip()
                    rectangles.append({
                        'text': text,
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height,
                    })
            except Exception:
                # If auto_shape_type check fails, still consider shapes with text
                text = ""
                if shape.has_text_frame:
                    text = " ".join(
                        p.text.strip() for p in shape.text_frame.paragraphs
                    ).strip()
                if text:
                    rectangles.append({
                        'text': text,
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height,
                    })
        # Check for connectors/lines
        if shape.shape_type in (MSO_SHAPE_TYPE.LINE, 9):
            connectors.append(shape)
        # Also check for freeform or other connector types
        if hasattr(shape, 'begin_x') and hasattr(shape, 'end_x'):
            if shape not in connectors:
                connectors.append(shape)

    expected_labels = ['Data Collection', 'Data Cleaning', 'Analysis', 'Visualization']

    # Component 1: Four rectangles with correct labels (0.5 points)
    # Award 0.125 per correct label found
    try:
        found_labels = [r['text'] for r in rectangles]
        matched_count = 0
        for label in expected_labels:
            if any(label.lower() in fl.lower() for fl in found_labels):
                matched_count += 1

        if matched_count == 4:
            print(f"PASS: Component 1 — All 4 rectangle labels found: {found_labels} (0.5 pts)")
            total_score += 0.5
        elif matched_count > 0:
            partial = matched_count * 0.125
            print(f"PARTIAL: Component 1 — {matched_count}/4 labels found: {found_labels} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — No matching rectangle labels found on slide 5. Found shapes with text: {found_labels}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Three connectors/arrows between rectangles (0.3 points)
    # Award 0.1 per connector found, up to 3
    try:
        num_connectors = len(connectors)
        if num_connectors >= 3:
            print(f"PASS: Component 2 — Found {num_connectors} connectors (0.3 pts)")
            total_score += 0.3
        elif num_connectors > 0:
            partial = min(num_connectors, 3) * 0.1
            print(f"PARTIAL: Component 2 — Found {num_connectors}/3 connectors ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No connectors/lines found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Horizontal arrangement in correct order (0.2 points)
    # Rectangles should be at roughly the same Y position and ordered left-to-right
    # matching: Data Collection -> Data Cleaning -> Analysis -> Visualization
    try:
        if len(rectangles) >= 4:
            # Match each expected label to its rectangle
            label_rects = {}
            for label in expected_labels:
                for r in rectangles:
                    if label.lower() in r['text'].lower():
                        label_rects[label] = r
                        break

            if len(label_rects) == 4:
                # Check horizontal alignment: all tops within 15% of slide height
                tops = [label_rects[l]['top'] for l in expected_labels]
                top_range = max(tops) - min(tops)
                slide_height = prs.slide_height
                horizontally_aligned = top_range < (slide_height * 0.15)

                # Check left-to-right order
                lefts = [label_rects[l]['left'] for l in expected_labels]
                correctly_ordered = all(lefts[i] < lefts[i + 1] for i in range(3))

                if horizontally_aligned and correctly_ordered:
                    print(f"PASS: Component 3 — Rectangles horizontally arranged in correct order (0.2 pts)")
                    total_score += 0.2
                elif horizontally_aligned:
                    print(f"PARTIAL: Component 3 — Horizontally aligned but wrong order. Lefts: {lefts} (0.1 pts)")
                    total_score += 0.1
                elif correctly_ordered:
                    print(f"PARTIAL: Component 3 — Correct order but not horizontally aligned. Top range: {top_range} (0.1 pts)")
                    total_score += 0.1
                else:
                    print(f"FAIL: Component 3 — Not horizontal (top range={top_range}) and not ordered (lefts={lefts})")
            else:
                print(f"FAIL: Component 3 — Could not match all 4 labels to rectangles. Matched: {list(label_rects.keys())}")
        else:
            print(f"FAIL: Component 3 — Need 4 rectangles for arrangement check, found {len(rectangles)}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
