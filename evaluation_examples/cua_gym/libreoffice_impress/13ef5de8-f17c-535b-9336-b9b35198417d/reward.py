"""
Reward Script: Project Status Report Presentation from Excel Data
Task ID: impress_wf_043
Domain: libreoffice_impress
Scoring:
  C1 (0.15): File exists with exactly 8 slides
  C2 (0.10): Slide 1 title contains "Project Phoenix"
  C3 (0.15): Slide 2 has 3 colored rectangles (status dashboard cards)
  C4 (0.10): Slide 3 has horizontal progress bar rectangles (milestone tracker)
  C5 (0.15): Slide 4 has a table with risk data and color-coded cells
  C6 (0.15): Slide 5 has a bar/column chart (resource utilization)
  C7 (0.15): Slide 6 has a line chart (sprint burndown)
  C8 (0.05): Slide 7 has blocker items with colored indicator rectangles
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_043'
FILE_PATH = os.path.join(WORKDIR, 'Desktop', 'Project_Status.pptx')


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: File has exactly 8 slides (0.15 points)
    try:
        num_slides = len(prs.slides)
        if num_slides == 8:
            print(f"PASS: Component 1 — 8 slides found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — expected 8 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 1 title contains "Project Phoenix" (0.10 points)
    try:
        slide1 = prs.slides[0]
        slide1_text = ""
        for shape in slide1.shapes:
            if shape.has_text_frame:
                slide1_text += " " + shape.text_frame.text
        if "project phoenix" in slide1_text.lower():
            print(f"PASS: Component 2 — Slide 1 contains 'Project Phoenix' (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 — 'Project Phoenix' not found in slide 1 text: {slide1_text[:100]}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 2 has 3 colored rectangles for status dashboard (0.15 points)
    # Task: 3 metric cards (Budget: On Track, Timeline: At Risk, Scope: On Track) with green/yellow/red backgrounds
    try:
        slide2 = prs.slides[1]
        colored_rects = []
        for shape in slide2.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                try:
                    fill = shape.fill
                    if fill.type is not None and fill.type == 1:  # SOLID fill
                        color = str(fill.fore_color.rgb).upper()
                        colored_rects.append(color)
                except Exception:
                    pass

        # Check slide 2 also has status text like "Budget", "Timeline", "Scope"
        slide2_text = ""
        for shape in slide2.shapes:
            if shape.has_text_frame:
                slide2_text += " " + shape.text_frame.text.lower()

        has_metric_labels = ("budget" in slide2_text and "timeline" in slide2_text and "scope" in slide2_text)
        has_status_text = ("on track" in slide2_text and "at risk" in slide2_text)

        if len(colored_rects) >= 3 and has_metric_labels and has_status_text:
            print(f"PASS: Component 3 — Slide 2 has {len(colored_rects)} colored rects with metric labels (0.15 pts)")
            total_score += 0.15
        elif len(colored_rects) >= 3 and has_metric_labels:
            print(f"PARTIAL: Component 3 — Rects and labels present but missing status text (0.10 pts)")
            total_score += 0.10
        elif len(colored_rects) >= 3:
            print(f"PARTIAL: Component 3 — Colored rects found but missing labels (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 3 — Found {len(colored_rects)} colored rects, need 3. Labels: {has_metric_labels}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 3 has horizontal progress bar rectangles (0.10 points)
    # Milestone tracker with rectangle pairs (background + progress bar)
    try:
        slide3 = prs.slides[2]
        rect_count = 0
        for shape in slide3.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                rect_count += 1

        # Should have multiple rectangle pairs (at least 4 rects for progress bars)
        slide3_text = ""
        for shape in slide3.shapes:
            if shape.has_text_frame:
                slide3_text += " " + shape.text_frame.text.lower()

        has_milestone_content = ("milestone" in slide3_text or "completed" in slide3_text or
                                  "pending" in slide3_text or "delayed" in slide3_text)

        if rect_count >= 4 and has_milestone_content:
            print(f"PASS: Component 4 — Slide 3 has {rect_count} rects with milestone content (0.10 pts)")
            total_score += 0.10
        elif rect_count >= 4:
            print(f"PARTIAL: Component 4 — {rect_count} rects but missing milestone text (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 4 — Only {rect_count} rectangles on slide 3, need >= 4")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slide 4 has a table with risk data and color-coded cells (0.15 points)
    try:
        slide4 = prs.slides[3]
        table_found = False
        has_risk_headers = False
        has_colored_text = False

        for shape in slide4.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_found = True
                tbl = shape.table
                rows = len(tbl.rows)
                cols = len(tbl.columns)

                # Check for risk-related headers in first row
                header_text = " ".join([tbl.cell(0, c).text.lower() for c in range(cols)])
                if "risk" in header_text and ("probability" in header_text or "impact" in header_text):
                    has_risk_headers = True

                # Check for color-coded text in cells (non-header rows)
                color_count = 0
                for r in range(1, min(rows, 7)):
                    for c in range(cols):
                        cell = tbl.cell(r, c)
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                try:
                                    if run.font.color.type is not None:
                                        rgb = str(run.font.color.rgb)
                                        if rgb != "000000":  # non-black = color-coded
                                            color_count += 1
                                except Exception:
                                    pass
                if color_count >= 2:
                    has_colored_text = True

        if table_found and has_risk_headers and has_colored_text:
            print(f"PASS: Component 5 — Risk table with color-coded cells found (0.15 pts)")
            total_score += 0.15
        elif table_found and has_risk_headers:
            print(f"PARTIAL: Component 5 — Risk table found but no color-coding (0.10 pts)")
            total_score += 0.10
        elif table_found:
            print(f"PARTIAL: Component 5 — Table found but missing risk headers (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 5 — No table found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Slide 5 has a bar/column chart (0.15 points)
    try:
        slide5 = prs.slides[4]
        chart_found = False
        is_bar_chart = False

        for shape in slide5.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                chart_found = True
                chart_type = shape.chart.chart_type
                # COLUMN_CLUSTERED=51, COLUMN_STACKED=52, BAR_CLUSTERED=57, BAR_STACKED=58
                if chart_type in (51, 52, 53, 54, 57, 58, 59, 60):
                    is_bar_chart = True
                print(f"  Chart type value: {chart_type}")

        if chart_found and is_bar_chart:
            print(f"PASS: Component 6 — Bar/column chart on slide 5 (0.15 pts)")
            total_score += 0.15
        elif chart_found:
            print(f"PARTIAL: Component 6 — Chart found but not bar/column type (0.07 pts)")
            total_score += 0.07
        else:
            print(f"FAIL: Component 6 — No chart found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Slide 6 has a line chart (0.15 points)
    try:
        slide6 = prs.slides[5]
        chart_found = False
        is_line_chart = False

        for shape in slide6.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                chart_found = True
                chart_type = shape.chart.chart_type
                # LINE=4, LINE_MARKERS=5, LINE_STACKED=6, etc.
                if chart_type in (4, 5, 6, 7, 8, 9):
                    is_line_chart = True
                print(f"  Chart type value: {chart_type}")

        if chart_found and is_line_chart:
            print(f"PASS: Component 7 — Line chart on slide 6 (0.15 pts)")
            total_score += 0.15
        elif chart_found:
            print(f"PARTIAL: Component 7 — Chart found but not line type (0.07 pts)")
            total_score += 0.07
        else:
            print(f"FAIL: Component 7 — No chart found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    # Component 8: Slide 7 has blocker items with colored indicator rectangles (0.05 points)
    try:
        slide7 = prs.slides[6]
        rect_count = 0
        for shape in slide7.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                rect_count += 1

        slide7_text = ""
        for shape in slide7.shapes:
            if shape.has_text_frame:
                slide7_text += " " + shape.text_frame.text.lower()

        has_blocker_content = ("blocker" in slide7_text or "priority" in slide7_text or
                               "critical" in slide7_text or "owner" in slide7_text)

        if rect_count >= 2 and has_blocker_content:
            print(f"PASS: Component 8 — Slide 7 has {rect_count} indicator rects with blocker content (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 8 — Rects: {rect_count}, blocker content: {has_blocker_content}")
    except Exception as e:
        print(f"ERROR: Component 8 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entrypoint
if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(FILE_PATH)
