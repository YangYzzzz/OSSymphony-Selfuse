"""
Initial Setup: Build a before/after comparison slide on slide 6
Task ID: impress_design_075
Domain: libreoffice_impress

Creates an 8-slide "Redesign_Case" presentation. Slide 6 has only a title
'Visual Comparison' — no divider, no labels, no placeholder rectangles.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_design_075'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.paragraphs[0].text = point
        else:
            p = tf.add_paragraph()
            p.text = point
    return slide


def add_blank_slide_with_title(prs, title_text):
    """Add a blank slide with just a title textbox at the top."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add a title textbox manually
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    return slide


def create_initial():
    prs = Presentation()
    # Standard widescreen 13.333 x 7.5 inches
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs, "Website Redesign Case Study",
                    "Transforming the Digital Experience for Meridian Health")

    # Slide 2: Project Overview
    add_content_slide(prs, "Project Overview", [
        "Client: Meridian Health Group - a network of 12 regional hospitals",
        "Objective: Modernize patient portal and public-facing website",
        "Timeline: January 2025 - June 2025 (6 months)",
        "Budget: $340,000 allocated across design, development, and QA",
        "Team: 4 designers, 6 developers, 2 QA engineers, 1 project manager",
    ])

    # Slide 3: Research Findings
    add_content_slide(prs, "Research Findings", [
        "72% of patients found appointment booking confusing",
        "Average task completion time: 4.2 minutes (industry avg: 1.8 min)",
        "Mobile traffic accounts for 68% of total visits",
        "Accessibility audit revealed 23 WCAG 2.1 AA violations",
        "Net Promoter Score dropped from 42 to 31 over 18 months",
    ])

    # Slide 4: Design Process
    add_content_slide(prs, "Design Process", [
        "Phase 1: Stakeholder interviews and user journey mapping",
        "Phase 2: Competitive analysis of 8 healthcare portals",
        "Phase 3: Wireframing with Figma - 47 screens prototyped",
        "Phase 4: Usability testing with 15 patients across 3 age groups",
        "Phase 5: Visual design system and component library creation",
    ])

    # Slide 5: Key Metrics
    add_content_slide(prs, "Key Metrics", [
        "Task completion rate improved from 58% to 91%",
        "Average booking time reduced to 1.4 minutes",
        "Mobile bounce rate decreased by 34%",
        "WCAG 2.1 AA compliance achieved (0 violations)",
        "Patient satisfaction score: 4.6/5.0 (up from 3.1/5.0)",
    ])

    # Slide 6: Visual Comparison - ONLY TITLE, nothing else
    add_blank_slide_with_title(prs, "Visual Comparison")

    # Slide 7: Implementation Timeline
    add_content_slide(prs, "Implementation Timeline", [
        "Q1 2025: Discovery, research, and information architecture",
        "Q2 2025: UI design, prototyping, and usability testing",
        "Q3 2025: Frontend and backend development sprints",
        "Q4 2025: QA testing, accessibility audit, and soft launch",
        "January 2026: Full production rollout across all 12 hospitals",
    ])

    # Slide 8: Next Steps
    add_content_slide(prs, "Next Steps", [
        "Integrate patient feedback loop for continuous improvement",
        "Expand telehealth scheduling to all specialty departments",
        "Develop native iOS and Android companion apps",
        "Launch multilingual support (Spanish, Mandarin, Vietnamese)",
        "Plan Phase 2: provider-facing dashboard redesign",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
