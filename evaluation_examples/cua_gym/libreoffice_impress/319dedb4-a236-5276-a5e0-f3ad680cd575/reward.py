"""
Reward Script: Before/After comparison slide on slide 6
Task ID: impress_design_075
Domain: libreoffice_impress
Scoring:
  Component 1: Vertical divider line at x~6.667in spanning slide height (0.20)
  Component 2: 'BEFORE' text label in 20pt bold #E74C3C at ~x=2in y=0.5in (0.20)
  Component 3: 'AFTER' text label in 20pt bold #27AE60 at ~x=8.5in y=0.5in (0.20)
  Component 4: Left placeholder rectangle ~5.5x5in at y~1.5in (0.20)
  Component 5: Right placeholder rectangle ~5.5x5in at y~1.5in (0.20)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_design_075'
FILE_NAME = 'impress_design_075.pptx'

# Tolerance for position/size comparisons (relative)
POS_TOLERANCE = 0.10  # 10% relative tolerance for positions
ABS_TOLERANCE = Inches(0.3)  # absolute tolerance for positions


def approx_eq(val, expected, rel_tol=POS_TOLERANCE, abs_tol=ABS_TOLERANCE):
    """Check if val is approximately equal to expected within tolerance."""
    if abs(val - expected) <= abs_tol:
        return True
    if expected == 0:
        return abs(val) <= abs_tol
    return abs(val - expected) / abs(expected) <= rel_tol


def find_shapes_by_text(slide, text_target):
    """Find shapes containing specific text (case-insensitive)."""
    results = []
    for shape in slide.shapes:
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            full_text = shape.text_frame.text.strip()
            if full_text.upper() == text_target.upper():
                results.append(shape)
    return results


def find_line_shapes(slide):
    """Find line/connector shapes on the slide."""
    lines = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            lines.append(shape)
        # Also check for freeform lines or connectors by XML tag
        elif 'cxnSp' in shape._element.tag:
            lines.append(shape)
    return lines


def find_rectangles(slide):
    """Find auto-shape rectangles on the slide."""
    rects = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # Check if it's a rectangle by name or by shape type
            if 'rectangle' in shape.name.lower() or 'rect' in shape.name.lower():
                rects.append(shape)
            else:
                # Try to check auto_shape_type
                try:
                    from pptx.enum.shapes import MSO_SHAPE
                    if shape.auto_shape_type == MSO_SHAPE.RECTANGLE:
                        rects.append(shape)
                except:
                    # If we can't determine, include it as candidate
                    rects.append(shape)
    return rects


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[5]  # slide 6, 0-indexed

    # Component 1: Vertical divider line at x~6.667in spanning slide height (0.20 pts)
    try:
        lines = find_line_shapes(slide)
        line_found = False
        expected_x = Inches(6.667)
        expected_height = prs.slide_height  # full slide height

        for line in lines:
            # Vertical line: width ~0, height ~slide height, x ~6.667in
            is_vertical = line.width <= Inches(0.1)  # effectively zero width
            is_at_x = approx_eq(line.left, expected_x)
            is_tall = approx_eq(line.height, expected_height, rel_tol=0.15)

            if is_vertical and is_at_x and is_tall:
                line_found = True
                print(f"PASS: Component 1 -- Vertical line at x={line.left/914400:.3f}in, "
                      f"height={line.height/914400:.3f}in (0.20 pts)")
                total_score += 0.20
                break

        if not line_found:
            if lines:
                for l in lines:
                    print(f"  Found line: left={l.left/914400:.3f}in, width={l.width/914400:.3f}in, "
                          f"height={l.height/914400:.3f}in")
            else:
                print("  No line shapes found on slide 6")
            print("FAIL: Component 1 -- No vertical divider line at x~6.667in")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: 'BEFORE' text in 20pt bold #E74C3C at x~2in y~0.5in (0.20 pts)
    try:
        before_shapes = find_shapes_by_text(slide, 'BEFORE')
        before_found = False

        for shape in before_shapes:
            # Check position
            x_ok = approx_eq(shape.left, Inches(2.0))
            y_ok = approx_eq(shape.top, Inches(0.5))

            if not (x_ok and y_ok):
                continue

            # Check font properties on the run
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip().upper() == 'BEFORE':
                        is_bold = run.font.bold is True
                        size_ok = run.font.size is not None and approx_eq(
                            run.font.size, Pt(20), rel_tol=0.05, abs_tol=Pt(1))
                        try:
                            color_ok = (run.font.color.type is not None and
                                        str(run.font.color.rgb).upper() == 'E74C3C')
                        except:
                            color_ok = False

                        if is_bold and size_ok and color_ok:
                            before_found = True
                            print(f"PASS: Component 2 -- 'BEFORE' at ({shape.left/914400:.2f}in, "
                                  f"{shape.top/914400:.2f}in), bold, 20pt, #E74C3C (0.20 pts)")
                            total_score += 0.20
                        else:
                            details = f"bold={is_bold}, size_ok={size_ok}, color_ok={color_ok}"
                            print(f"PARTIAL FAIL: Component 2 -- 'BEFORE' found but properties wrong: {details}")
                            # Give partial credit for having the text in roughly right position
                            if x_ok and y_ok:
                                total_score += 0.05
                                print(f"  Partial credit: 0.05 for text placement")
                            before_found = True
                        break
            if before_found:
                break

        if not before_found:
            print(f"FAIL: Component 2 -- 'BEFORE' text not found at expected position")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: 'AFTER' text in 20pt bold #27AE60 at x~8.5in y~0.5in (0.20 pts)
    try:
        after_shapes = find_shapes_by_text(slide, 'AFTER')
        after_found = False

        for shape in after_shapes:
            x_ok = approx_eq(shape.left, Inches(8.5))
            y_ok = approx_eq(shape.top, Inches(0.5))

            if not (x_ok and y_ok):
                continue

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip().upper() == 'AFTER':
                        is_bold = run.font.bold is True
                        size_ok = run.font.size is not None and approx_eq(
                            run.font.size, Pt(20), rel_tol=0.05, abs_tol=Pt(1))
                        try:
                            color_ok = (run.font.color.type is not None and
                                        str(run.font.color.rgb).upper() == '27AE60')
                        except:
                            color_ok = False

                        if is_bold and size_ok and color_ok:
                            after_found = True
                            print(f"PASS: Component 3 -- 'AFTER' at ({shape.left/914400:.2f}in, "
                                  f"{shape.top/914400:.2f}in), bold, 20pt, #27AE60 (0.20 pts)")
                            total_score += 0.20
                        else:
                            details = f"bold={is_bold}, size_ok={size_ok}, color_ok={color_ok}"
                            print(f"PARTIAL FAIL: Component 3 -- 'AFTER' found but properties wrong: {details}")
                            if x_ok and y_ok:
                                total_score += 0.05
                                print(f"  Partial credit: 0.05 for text placement")
                            after_found = True
                        break
            if after_found:
                break

        if not after_found:
            print(f"FAIL: Component 3 -- 'AFTER' text not found at expected position")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4 & 5: Two placeholder rectangles (~5.5x5in) at y~1.5in (0.20 pts each)
    try:
        rects = find_rectangles(slide)
        expected_w = Inches(5.5)
        expected_h = Inches(5.0)
        expected_y = Inches(1.5)

        # Separate into left-half and right-half rectangles
        slide_midpoint = prs.slide_width / 2
        left_rect_found = False
        right_rect_found = False

        for rect in rects:
            w_ok = approx_eq(rect.width, expected_w, rel_tol=0.15)
            h_ok = approx_eq(rect.height, expected_h, rel_tol=0.15)
            y_ok = approx_eq(rect.top, expected_y)

            if w_ok and h_ok and y_ok:
                center_x = rect.left + rect.width / 2
                if center_x < slide_midpoint and not left_rect_found:
                    left_rect_found = True
                    print(f"PASS: Component 4 -- Left rectangle at ({rect.left/914400:.2f}in, "
                          f"{rect.top/914400:.2f}in), size {rect.width/914400:.2f}x{rect.height/914400:.2f}in (0.20 pts)")
                    total_score += 0.20
                elif center_x >= slide_midpoint and not right_rect_found:
                    right_rect_found = True
                    print(f"PASS: Component 5 -- Right rectangle at ({rect.left/914400:.2f}in, "
                          f"{rect.top/914400:.2f}in), size {rect.width/914400:.2f}x{rect.height/914400:.2f}in (0.20 pts)")
                    total_score += 0.20

        if not left_rect_found:
            print(f"FAIL: Component 4 -- No left placeholder rectangle found")
            if rects:
                for r in rects:
                    print(f"  Candidate rect: left={r.left/914400:.2f}in, top={r.top/914400:.2f}in, "
                          f"w={r.width/914400:.2f}in, h={r.height/914400:.2f}in")

        if not right_rect_found:
            print(f"FAIL: Component 5 -- No right placeholder rectangle found")

    except Exception as e:
        print(f"ERROR: Components 4/5 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{FILE_NAME}'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
