"""
Reward Script: Wedding presentation slideshow verification
Task ID: impress_wf_051
Domain: libreoffice_impress
Scoring:
  Component 1: File exists on Desktop (0.10)
  Component 2: Exactly 10 slides (0.10)
  Component 3: Background color #F8BBD0 on all slides (0.15)
  Component 4: Text color #880E4F (burgundy) on main text (0.10)
  Component 5: Slide 1 has italic title text + decorative line shapes (0.10)
  Component 6: Slide 2 has 3 timeline events (0.10)
  Component 7: Slides 3-7 have rectangle placeholders + italic captions (0.10)
  Component 8: Slide 8 has a heart shape (0.10)
  Component 9: Dissolve transitions with 6-second auto-advance (0.15)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_051'
FILE_PATH = os.path.join(WORKDIR, 'Desktop', 'Wedding_Story.pptx')


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
    from pptx.dml.color import RGBColor

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: File exists on Desktop (0.10 points)
    # This is implicitly passed since we loaded it, but we check the path
    # NOTE: File existence alone is not scored — we already loaded it.
    # We check it's specifically on the Desktop (task requirement: "Save as ... on the Desktop")
    try:
        if os.path.dirname(os.path.abspath(file_path)) == os.path.join(WORKDIR, 'Desktop'):
            print(f"PASS: Component 1 — File is on Desktop (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 1 — File not on Desktop, found at {file_path}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Exactly 10 slides (0.10 points)
    try:
        slide_count = len(prs.slides)
        if slide_count == 10:
            print(f"PASS: Component 2 — 10 slides found (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 — Expected 10 slides, found {slide_count}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Background color #F8BBD0 on all slides (0.15 points)
    try:
        bg_pass_count = 0
        for idx, slide in enumerate(prs.slides):
            fill = slide.background.fill
            if fill.type is not None:
                try:
                    rgb = str(fill.fore_color.rgb).upper()
                    if rgb == 'F8BBD0':
                        bg_pass_count += 1
                    else:
                        print(f"  Slide {idx+1} background: {rgb} (expected F8BBD0)")
                except Exception:
                    print(f"  Slide {idx+1} background: cannot read RGB")
            else:
                print(f"  Slide {idx+1} background: inherited/none")

        if bg_pass_count == len(prs.slides) and len(prs.slides) > 0:
            print(f"PASS: Component 3 — All {bg_pass_count} slides have #F8BBD0 background (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 3 — {bg_pass_count}/{len(prs.slides)} slides have correct background")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Text color #880E4F (burgundy) on main text (0.10 points)
    # Check that at least 80% of non-empty text runs across all slides use burgundy color
    try:
        burgundy_runs = 0
        total_runs = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if (run.text or "").strip():
                                total_runs += 1
                                try:
                                    if run.font.color.type is not None:
                                        rgb = str(run.font.color.rgb).upper()
                                        if rgb == '880E4F':
                                            burgundy_runs += 1
                                except Exception:
                                    pass

        if total_runs > 0 and burgundy_runs / total_runs >= 0.7:
            print(f"PASS: Component 4 — {burgundy_runs}/{total_runs} text runs use #880E4F (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 4 — Only {burgundy_runs}/{total_runs} text runs use #880E4F")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slide 1 has italic title + decorative line shapes above/below (0.10 points)
    try:
        slide1 = prs.slides[0]
        comp5_score = 0.0

        # Check for italic text containing "Sarah & Michael"
        has_italic_title = False
        for shape in slide1.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    full_text = para.text.strip()
                    if 'sarah' in full_text.lower() and 'michael' in full_text.lower():
                        for run in para.runs:
                            if run.font.italic is True:
                                has_italic_title = True
                                break

        # Check for rectangular line shapes (thin rectangles used as decorative lines)
        rect_count = 0
        for shape in slide1.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    # Thin rectangles (height < 0.5 inch = ~457200 EMU) are decorative lines
                    if shape.height < 457200:
                        rect_count += 1
            except Exception:
                pass

        if has_italic_title and rect_count >= 2:
            print(f"PASS: Component 5 — Italic title with {rect_count} decorative lines (0.10 pts)")
            comp5_score = 0.10
        elif has_italic_title:
            print(f"PARTIAL: Component 5 — Italic title found but only {rect_count} decorative line(s) (0.05 pts)")
            comp5_score = 0.05
        else:
            print(f"FAIL: Component 5 — italic_title={has_italic_title}, line_shapes={rect_count}")

        total_score += comp5_score
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Slide 2 has 3 timeline events (0.10 points)
    try:
        slide2 = prs.slides[1]

        # Count text boxes that appear to be date/event entries
        # Timeline events have bold date text boxes
        date_boxes = 0
        has_title = False
        for shape in slide2.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    if 'how we met' in text.lower():
                        has_title = True
                    # Date entries typically contain year numbers and are bold
                    for run in para.runs:
                        if run.font.bold is True and any(str(yr) in text for yr in range(2000, 2030)):
                            date_boxes += 1
                            break

        if date_boxes >= 3:
            print(f"PASS: Component 6 — {date_boxes} timeline date entries found (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 6 — Expected 3+ timeline events, found {date_boxes} date entries")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Slides 3-7 have rectangle placeholders + italic captions (0.10 points)
    try:
        slides_with_rect_and_italic = 0
        for idx in range(2, 7):  # slides 3-7 (0-indexed: 2-6)
            if idx >= len(prs.slides):
                continue
            slide = prs.slides[idx]

            has_rect = False
            has_italic_caption = False
            for shape in slide.shapes:
                # Check for rectangle placeholder
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        has_rect = True
                except Exception:
                    pass

                # Check for italic caption text
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text and 'photo' not in text.lower():
                            for run in para.runs:
                                if run.font.italic is True:
                                    has_italic_caption = True
                                    break

            if has_rect and has_italic_caption:
                slides_with_rect_and_italic += 1

        if slides_with_rect_and_italic == 5:
            print(f"PASS: Component 7 — All 5 photo slides (3-7) have rectangles + italic captions (0.10 pts)")
            total_score += 0.10
        elif slides_with_rect_and_italic >= 3:
            partial = round(0.10 * slides_with_rect_and_italic / 5, 2)
            print(f"PARTIAL: Component 7 — {slides_with_rect_and_italic}/5 slides have rect+italic ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 7 — Only {slides_with_rect_and_italic}/5 photo slides have rect+italic caption")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    # Component 8: Slide 8 has a heart shape (0.10 points)
    try:
        slide8 = prs.slides[7]
        has_heart = False
        for shape in slide8.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    if shape.auto_shape_type == MSO_SHAPE.HEART:
                        has_heart = True
                        break
            except Exception:
                pass

        if has_heart:
            print(f"PASS: Component 8 — Heart shape found on slide 8 (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 8 — No heart shape found on slide 8")
    except Exception as e:
        print(f"ERROR: Component 8 — {e}")

    # Component 9: Dissolve transitions with 6-second auto-advance on all slides (0.15 points)
    try:
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        dissolve_count = 0
        auto_advance_count = 0
        slide_total = len(prs.slides)

        with zipfile.ZipFile(file_path, 'r') as zf:
            for i in range(1, slide_total + 1):
                try:
                    with zf.open(f'ppt/slides/slide{i}.xml') as f:
                        root = ET.parse(f).getroot()
                        tr = root.find('.//p:transition', ns)
                        if tr is not None:
                            # Check dissolve
                            if tr.find('.//p:dissolve', ns) is not None:
                                dissolve_count += 1
                            # Check auto-advance time (advTm in milliseconds)
                            adv_tm = tr.get('advTm')
                            if adv_tm is not None and int(adv_tm) == 6000:
                                auto_advance_count += 1
                except Exception:
                    pass

        if dissolve_count == slide_total and auto_advance_count == slide_total:
            print(f"PASS: Component 9 — All {slide_total} slides have dissolve + 6s advance (0.15 pts)")
            total_score += 0.15
        elif dissolve_count > 0 or auto_advance_count > 0:
            partial = round(0.15 * (dissolve_count + auto_advance_count) / (2 * slide_total), 2)
            print(f"PARTIAL: Component 9 — dissolve={dissolve_count}/{slide_total}, advance={auto_advance_count}/{slide_total} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 9 — No dissolve transitions or auto-advance found")
    except Exception as e:
        print(f"ERROR: Component 9 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point: test against canonical artifact path
if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(FILE_PATH)
