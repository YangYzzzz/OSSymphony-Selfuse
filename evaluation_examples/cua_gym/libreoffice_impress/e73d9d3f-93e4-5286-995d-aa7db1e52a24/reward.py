"""
FINAL REWARD SCRIPT - SUCCESS
Task: Add page numbers to the footer on pages 2–5 only, centered.
Generated: 2025-10-17 09:03:42
Status: success
Model: azure-o3
Total Steps: 8
"""

import os
import re
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER

def _detect_page_number_shape(slide):
    """Detect if the slide contains a page-number shape.

    Returns
    -------
    (found, centered)
        found    : bool – True if a page-number placeholder/text is detected
        centered : bool – True if the detected shape is center-aligned
    """
    for shape in slide.shapes:
        # 1. Check for an explicit slide-number placeholder
        try:
            if shape.is_placeholder and \
               shape.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
                centered = False
                if shape.has_text_frame and shape.text_frame.paragraphs:
                    centered = (
                        shape.text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER
                    )
                return True, centered
        except Exception:
            pass  # Not a valid placeholder – safely ignore

        # 2. Fallback: look for a text frame that is either
        #    the placeholder token (“<#>”) or just a digit
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt == "<#>" or re.fullmatch(r"\d+", txt):
                centered = (
                    shape.text_frame.paragraphs and
                    shape.text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER
                )
                return True, centered
    # Nothing that looks like a page number was found
    return False, False

def verify_task(file_path):
    """Reward-function for the task:
    "Add page numbers to the footer on pages 2–5 only, centered."""

    # ----- Basic file validation (NO POINTS AWARDED) -----
    if not os.path.exists(file_path):
        print("✗ File not found:", file_path)
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print("✗ Unable to open PPTX:", e)
        return 0.0

    total_slides = len(prs.slides)
    print(f"Loaded presentation with {total_slides} slides\n")

    # Slides that MUST have page numbers
    required_slides = list(range(2, 6))  # 2,3,4,5 (1-based indexing)

    # ---------------- Scoring weights ----------------
    per_slide_presence   = 0.15   # 0.15 × 4  = 0.60
    per_slide_alignment  = 0.05   # 0.05 × 4  = 0.20
    absence_full_score   = 0.20   #            = 0.20
    # -------------------------------------------------

    presence_score  = 0.0
    alignment_score = 0.0

    # ---------- 1. Check required slides ----------
    for idx in required_slides:
        if idx > total_slides:
            print(f"Slide {idx} missing – cannot verify requirements.")
            continue
        found, centered = _detect_page_number_shape(prs.slides[idx-1])
        print(
            f"Slide {idx}: page-number found={found}, center-aligned={centered}"
        )

        if found:
            presence_score += per_slide_presence
        if centered:
            alignment_score += per_slide_alignment
    print()

    # ---------- 2. Ensure absence on all other slides ----------
    absence_ok = True
    for idx in range(1, total_slides + 1):
        if idx in required_slides:
            continue  # already evaluated
        found, _ = _detect_page_number_shape(prs.slides[idx-1])
        print(f"Other slide {idx}: page-number present={found}")
        if found:
            absence_ok = False
    print()

    absence_score = absence_full_score if absence_ok else 0.0

    # ---------------- Final score ----------------
    total_score = presence_score + alignment_score + absence_score
    total_score = min(total_score, 1.0)

    print("Score breakdown:")
    print(f"  Presence (2-5):   {presence_score:.2f} / 0.60")
    print(f"  Alignment (2-5):  {alignment_score:.2f} / 0.20")
    print(f"  Absence (others): {absence_score:.2f} / 0.20")
    print(f"TOTAL SCORE: {total_score}")
    print(f"REWARD: {total_score}")

    return total_score

# ------------- Execute verification -------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/add_page_numbers_to_the_footer_on_pages_25_only_centered.pptx"
    verify_task(FILE_PATH)

