"""
FINAL REWARD SCRIPT - SUCCESS
Task: I’m up to slide 144 and just noticed the entire deck is still in the default 16:9 format. In LibreOffice Impress, how can I flip every slide in the file over to a 4:3 page size all at once without messing up the content?
Generated: 2025-09-10 18:07:02
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation

def verify_convert_to_4x3(file_path: str) -> float:
    """Reward script for the LibreOffice Impress task.

    The task: make sure the deck has been switched from 16:9 to 4:3 *and* that
    slide content still fits.  This script awards up to 1.0 points:
        • 0.6 pts – slide size matches a 4:3 aspect-ratio (±5 % tolerance)
        • 0.4 pts – all measurable shapes remain inside the new slide bounds
          (≤ 0 % overflow → full 0.4; ≤ 5 % overflow → 0.2; otherwise 0.0)

    It prints detailed diagnostics and always outputs the score as
    "REWARD: X.X".  A score of 1.0 means perfect completion.
    """

    print(f"Starting verification for: {file_path}\n")

    # ---------- Basic existence / load (no points) ----------
    if not os.path.exists(file_path):
        print("✗ File does not exist. Verification failed.")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Unable to load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # ---------- Requirement 1: page size is 4:3 ----------
    width, height = prs.slide_width, prs.slide_height
    ratio = width / height if height else 0
    target_ratio = 4 / 3
    size_score = 0.0

    print(f"Slide dimensions (EMU): width={width}, height={height}, ratio≈{ratio:.3f}")

    if abs(ratio - target_ratio) < 0.05:
        print("✓ Slide size is very close to 4:3 ratio")
        size_score = 0.6
    elif abs(ratio - target_ratio) < 0.10:
        print("⚠ Slide ratio slightly off 4:3 – giving partial credit (0.3)")
        size_score = 0.3
    else:
        print("✗ Slide ratio does not match 4:3 – no credit for size requirement")

    # ---------- Requirement 2: content still fits ----------
    content_score = 0.0
    total_shapes, overflowing_shapes = 0, 0

    for slide in prs.slides:
        for shape in slide.shapes:
            # Ensure the shape has geometry attributes
            if not all(hasattr(shape, a) for a in ("left", "top", "width", "height")):
                continue

            total_shapes += 1

            try:
                right = shape.left + shape.width
                bottom = shape.top + shape.height
            except Exception:
                # If reading geometry fails, skip shape (no credit either way)
                continue

            # Allow 2 % tolerance around edges
            if (right > width * 1.02 or bottom > height * 1.02 or
                    shape.left < -width * 0.02 or shape.top < -height * 0.02):
                overflowing_shapes += 1

    if total_shapes == 0:
        print("⚠ No measurable shapes found – cannot verify content placement. No points awarded for content fit.")
    else:
        overflow_ratio = overflowing_shapes / total_shapes
        print(f"Total measurable shapes: {total_shapes}")
        print(f"Shapes overflowing slide bounds: {overflowing_shapes} ({overflow_ratio:.2%})")

        if overflowing_shapes == 0:
            print("✓ All shapes fit within the slide bounds – full content credit (0.4)")
            content_score = 0.4
        elif overflow_ratio <= 0.05:
            print("⚠ Small number of shapes overflow – partial content credit (0.2)")
            content_score = 0.2
        else:
            print("✗ Significant content overflow – no content credit")

    # ---------- Final score ----------
    total_score = round(min(size_score + content_score, 1.0), 2)

    print("\nVerification complete.")
    print(f"REWARD: {total_score}")
    return total_score

# When run directly, verify the provided presentation path
if __name__ == "__main__":
    FILE_PATH = "/home/user/im_up_to_slide_144_and_just_noticed_the_entire_deck_is_still_in_the_default_169_format_in_libreoffic_golden.pptx"
    verify_convert_to_4x3(FILE_PATH)
