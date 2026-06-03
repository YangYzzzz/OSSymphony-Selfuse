"""
Initial Setup: Sales pitch presentation with speaker notes for extraction task
Task ID: osworld_multi_apps_impress_notes_export_007
Domain: libreoffice_impress (multi-app: Impress + Writer)

Creates SalesDeck.pptx on the Desktop with 9 slides of B2B sales pitch content
and multi-paragraph speaker notes on several slides.
The agent must extract notes and save them to sales_notes.docx on the Desktop.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_impress_notes_export_007'
OUTPUT = f'{WORKDIR}/SalesDeck.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_notes(slide, notes_text):
    """Set speaker notes on a slide. notes_text is a single string with \n for paragraph breaks."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    # Clear existing paragraphs
    tf.clear()
    paragraphs = notes_text.split('\n')
    for i, para_text in enumerate(paragraphs):
        if i == 0:
            tf.paragraphs[0].text = para_text
        else:
            p = tf.add_paragraph()
            p.text = para_text


def create_initial():
    prs = Presentation()
    # Standard widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]   # Title Slide
    content_layout = prs.slide_layouts[1]        # Title + Content
    title_only_layout = prs.slide_layouts[5]     # Blank (or title only)

    # ─── Slide 1: Title Slide ───────────────────────────────────────────────
    slide1 = prs.slides.add_slide(title_slide_layout)
    slide1.shapes.title.text = "TechVision Solutions"
    slide1.placeholders[1].text = "Transforming Enterprise Operations\nB2B Sales Presentation 2025"

    set_notes(slide1,
        "Welcome everyone to today's presentation.\n"
        "TechVision Solutions has been serving enterprise clients since 2015.\n"
        "Today we will walk through our full product suite and value proposition."
    )

    # ─── Slide 2: Agenda ────────────────────────────────────────────────────
    slide2 = prs.slides.add_slide(content_layout)
    slide2.shapes.title.text = "Agenda"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Company Overview"
    items2 = [
        "Market Challenges",
        "Our Solution Portfolio",
        "Client Success Stories",
        "Pricing & Packages",
        "Implementation Timeline",
        "Next Steps",
    ]
    for item in items2:
        p = tf2.add_paragraph()
        p.text = item
        p.level = 1

    set_notes(slide2,
        "Keep this slide brief — spend no more than 2 minutes on the agenda.\n"
        "Emphasize that we will have a Q&A session at the end.\n"
        "Ask if attendees have any specific areas of interest to prioritize."
    )

    # ─── Slide 3: Company Overview ──────────────────────────────────────────
    slide3 = prs.slides.add_slide(content_layout)
    slide3.shapes.title.text = "Company Overview"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Founded: 2015 | Headquarters: San Francisco, CA"
    stats = [
        "500+ enterprise clients across 30 countries",
        "97% client retention rate over the past 3 years",
        "ISO 27001 certified — enterprise-grade security",
        "24/7 dedicated support with SLA guarantees",
        "$120M ARR — Series C funded",
    ]
    for stat in stats:
        p = tf3.add_paragraph()
        p.text = stat
        p.level = 1

    set_notes(slide3,
        "Highlight our growth trajectory — 3x revenue in 2 years.\n"
        "Mention key customers in their industry vertical if known beforehand.\n"
        "The ISO certification often addresses security objections early.\n"
        "Reference the Series C as a signal of market confidence and stability."
    )

    # ─── Slide 4: Market Challenges ─────────────────────────────────────────
    slide4 = prs.slides.add_slide(content_layout)
    slide4.shapes.title.text = "The Challenges You Face"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Today's enterprise landscape demands:"
    challenges = [
        "Fragmented data silos blocking real-time decisions",
        "Manual workflows consuming 40% of operational capacity",
        "Rising cybersecurity risks and compliance overhead",
        "Difficulty scaling without proportional headcount growth",
        "Customer experience gaps costing 15-20% annual churn",
    ]
    for c in challenges:
        p = tf4.add_paragraph()
        p.text = c
        p.level = 1

    set_notes(slide4,
        "Pause here and ask: 'Which of these resonate most with your team?'\n"
        "Listen carefully — their answer will shape the rest of the conversation.\n"
        "If they mention data silos, pivot heavily to our DataBridge module.\n"
        "If compliance comes up, focus on our AuditTrail and policy engine features."
    )

    # ─── Slide 5: Solution Portfolio ────────────────────────────────────────
    slide5 = prs.slides.add_slide(content_layout)
    slide5.shapes.title.text = "Our Solution Portfolio"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Three integrated platforms, one unified experience:"
    solutions = [
        "DataBridge — real-time data integration and analytics",
        "FlowEngine — intelligent workflow automation",
        "SecureVault — compliance and risk management",
    ]
    for s in solutions:
        p = tf5.add_paragraph()
        p.text = s
        p.level = 1

    set_notes(slide5,
        "All three platforms integrate natively — no middleware required.\n"
        "Clients typically start with one module and expand within 6 months.\n"
        "DataBridge is our most popular entry point for data-heavy industries."
    )

    # ─── Slide 6: Client Success Stories ────────────────────────────────────
    slide6 = prs.slides.add_slide(content_layout)
    slide6.shapes.title.text = "Client Success Stories"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Results our clients have achieved:"
    stories = [
        "NexaCorp (Manufacturing): 35% reduction in supply chain delays",
        "FinEdge Bank: Full SOX compliance achieved in 90 days",
        "RetailPlus Group: 28% increase in customer retention",
        "HealthCare United: $4.2M annual savings from automation",
    ]
    for story in stories:
        p = tf6.add_paragraph()
        p.text = story
        p.level = 1

    set_notes(slide6,
        "Choose the 1-2 case studies most relevant to this prospect's industry.\n"
        "For financial services clients, lead with FinEdge Bank's compliance story.\n"
        "For manufacturing clients, NexaCorp's supply chain results are most compelling.\n"
        "Offer to share the full case study PDF after the meeting — it is a good leave-behind.\n"
        "Mention ROI timeline: most clients see positive ROI within 6-9 months."
    )

    # ─── Slide 7: Pricing & Packages ────────────────────────────────────────
    slide7 = prs.slides.add_slide(content_layout)
    slide7.shapes.title.text = "Pricing & Packages"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "Flexible pricing tailored to your scale:"
    tiers = [
        "Starter — up to 50 users — $3,500/month",
        "Business — up to 250 users — $9,200/month",
        "Enterprise — unlimited users — custom pricing",
        "All tiers include onboarding, training, and 24/7 support",
    ]
    for tier in tiers:
        p = tf7.add_paragraph()
        p.text = tier
        p.level = 1

    set_notes(slide7,
        "Do not lead with price — always establish value first.\n"
        "If the prospect pushes back on cost, redirect to total cost of ownership.\n"
        "The Business tier is our most common starting point for mid-market accounts.\n"
        "Enterprise deals are negotiated separately with volume and multi-year discounts."
    )

    # ─── Slide 8: Implementation Timeline ───────────────────────────────────
    slide8 = prs.slides.add_slide(content_layout)
    slide8.shapes.title.text = "Implementation Timeline"
    tf8 = slide8.placeholders[1].text_frame
    tf8.text = "From signed contract to go-live in 6 weeks:"
    timeline = [
        "Week 1-2: Discovery, kickoff, and environment setup",
        "Week 3-4: Data migration, integration configuration",
        "Week 5: User acceptance testing and training",
        "Week 6: Go-live and hypercare support",
        "Ongoing: Quarterly business reviews and optimization",
    ]
    for t in timeline:
        p = tf8.add_paragraph()
        p.text = t
        p.level = 1

    set_notes(slide8,
        "Emphasize our structured onboarding process — it reduces risk significantly.\n"
        "Our dedicated implementation team has a 100% on-time go-live record.\n"
        "The hypercare period in week 6 includes daily check-ins with our CSM team."
    )

    # ─── Slide 9: Next Steps ─────────────────────────────────────────────────
    slide9 = prs.slides.add_slide(content_layout)
    slide9.shapes.title.text = "Next Steps"
    tf9 = slide9.placeholders[1].text_frame
    tf9.text = "Ready to move forward?"
    steps = [
        "Schedule a technical deep-dive with your IT team",
        "Request a personalized demo with your data",
        "Review and sign the Master Service Agreement",
        "Kick off with your dedicated Customer Success Manager",
        "Contact: sales@techvision.io | +1 (415) 822-9100",
    ]
    for step in steps:
        p = tf9.add_paragraph()
        p.text = step
        p.level = 1

    set_notes(slide9,
        "Ask for a commitment on the next step before leaving the room.\n"
        "Ideal outcome: scheduled follow-up within 5 business days.\n"
        "Leave business cards and a printed one-pager summary.\n"
        "Follow up with a personalized email within 24 hours referencing key discussion points."
    )

    # Ensure Desktop directory exists
    os.makedirs(WORKDIR, exist_ok=True)
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open SalesDeck.pptx in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
