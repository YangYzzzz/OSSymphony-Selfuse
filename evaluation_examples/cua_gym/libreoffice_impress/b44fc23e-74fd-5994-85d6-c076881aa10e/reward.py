"""
Reward Script: Extract presenter notes from SalesDeck.pptx to sales_notes.docx
Task ID: osworld_multi_apps_impress_notes_export_007
Domain: libreoffice_impress (multi-app: impress + writer)
Scoring:
  Component 1 (0.25): sales_notes.docx exists on Desktop
  Component 2 (0.40): All 9 slides' notes text content is present in the docx
  Component 3 (0.20): Blank-line separators exist between slide note groups
  Component 4 (0.15): Internal paragraph breaks within slides are preserved
"""

import os

WORKDIR = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_impress_notes_export_007'

DOCX_PATH = f'{WORKDIR}/sales_notes.docx'
PPTX_PATH = f'{WORKDIR}/SalesDeck.pptx'


def get_slide_notes_paragraphs(slide):
    """Return a list of non-empty paragraph texts from a slide's notes."""
    try:
        ntf = slide.notes_slide.notes_text_frame
        paras = []
        for para in ntf.paragraphs:
            text = para.text.strip()
            if text:
                paras.append(text)
        return paras
    except Exception:
        return []


def verify_task():
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Component 1: sales_notes.docx exists on Desktop (0.25 points)
    # This FAILS on initial_env (no docx) and PASSES on golden_env
    try:
        if os.path.exists(DOCX_PATH):
            print(f"PASS: Component 1 — sales_notes.docx exists at {DOCX_PATH} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — sales_notes.docx not found at {DOCX_PATH}")
            print(f"\nScore: {total_score}/1.0")
            print(f"REWARD: {total_score}")
            return total_score
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Load the source PPTX to extract expected notes content
    expected_notes_per_slide = []
    try:
        from pptx import Presentation
        prs = Presentation(PPTX_PATH)
        num_slides = len(prs.slides)
        for slide in prs.slides:
            paras = get_slide_notes_paragraphs(slide)
            expected_notes_per_slide.append(paras)
        print(f"INFO: SalesDeck.pptx loaded — {num_slides} slides")
    except Exception as e:
        print(f"ERROR: Cannot load SalesDeck.pptx: {e}")
        # Cannot proceed with content checks; give partial credit for file existence
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Load the output docx
    try:
        from docx import Document
        doc = Document(DOCX_PATH)
        doc_paragraphs = [para.text for para in doc.paragraphs]
        # Collect all non-empty text lines in docx
        doc_nonempty_lines = [p.strip() for p in doc_paragraphs if p.strip()]
    except Exception as e:
        print(f"CRITICAL: Cannot load sales_notes.docx: {e}")
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: All slide notes text content is present in the docx (0.40 points)
    # Check representative lines from each slide's notes are present in the docx.
    # Collect the first paragraph text from each slide that has notes.
    try:
        slides_with_notes = [(i, paras) for i, paras in enumerate(expected_notes_per_slide) if paras]
        found_slides = 0
        total_slides_with_notes = len(slides_with_notes)

        for slide_idx, paras in slides_with_notes:
            # Check that first paragraph from this slide's notes appears in the docx
            first_para = paras[0].strip()
            found_in_doc = any(first_para in line or line in first_para for line in doc_nonempty_lines)
            if found_in_doc:
                found_slides += 1
            else:
                print(f"FAIL: Component 2 — Slide {slide_idx+1} first note para not found in docx: {repr(first_para[:80])}")

        if total_slides_with_notes > 0:
            coverage_ratio = found_slides / total_slides_with_notes
            component2_score = round(0.40 * coverage_ratio, 4)
            if coverage_ratio == 1.0:
                print(f"PASS: Component 2 — All {total_slides_with_notes} slides' notes found in docx (0.40 pts)")
                total_score += component2_score
            elif coverage_ratio > 0:
                print(f"PARTIAL: Component 2 — {found_slides}/{total_slides_with_notes} slides' notes found ({component2_score:.2f} pts)")
                total_score += component2_score
            else:
                print(f"FAIL: Component 2 — No slide notes found in docx")
        else:
            print("INFO: Component 2 — No slides with notes found in PPTX to verify")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Blank-line separators between slide notes groups (0.20 points)
    # The docx should have empty paragraphs separating notes from different slides.
    # There should be at least (number_of_slides_with_notes - 1) empty paragraph boundaries.
    try:
        # Count empty paragraphs in the docx (separators between groups)
        empty_para_count = sum(1 for p in doc_paragraphs if p.strip() == '')

        # We expect at least (slides_with_notes - 1) blank line separators
        expected_separators = max(0, total_slides_with_notes - 1)

        if expected_separators == 0:
            print("INFO: Component 3 — Only one slide with notes; no separator needed")
            total_score += 0.20
        elif empty_para_count >= expected_separators:
            print(f"PASS: Component 3 — Found {empty_para_count} blank-line separators (expected >= {expected_separators}) (0.20 pts)")
            total_score += 0.20
        elif empty_para_count > 0:
            partial_sep = round(0.20 * (empty_para_count / expected_separators), 4)
            print(f"PARTIAL: Component 3 — Found {empty_para_count} separators, expected {expected_separators} ({partial_sep:.2f} pts)")
            total_score += partial_sep
        else:
            print(f"FAIL: Component 3 — No blank-line separators found; expected {expected_separators}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Internal paragraph breaks within slides are preserved (0.15 points)
    # For slides that have multiple paragraphs of notes, verify that the individual
    # paragraph texts (beyond the first) also appear in the docx.
    try:
        multi_para_slides = [(i, paras) for i, paras in enumerate(expected_notes_per_slide) if len(paras) > 1]
        if not multi_para_slides:
            print("INFO: Component 4 — No multi-paragraph slides to verify")
            total_score += 0.15
        else:
            found_additional = 0
            total_additional = 0
            for slide_idx, paras in multi_para_slides:
                # Check paragraphs beyond the first
                for para_text in paras[1:]:
                    total_additional += 1
                    para_stripped = para_text.strip()
                    found_in_doc = any(para_stripped in line or line in para_stripped for line in doc_nonempty_lines)
                    if found_in_doc:
                        found_additional += 1
                    else:
                        print(f"FAIL: Component 4 — Slide {slide_idx+1} additional para not found: {repr(para_stripped[:80])}")

            if total_additional > 0:
                coverage = found_additional / total_additional
                comp4_score = round(0.15 * coverage, 4)
                if coverage == 1.0:
                    print(f"PASS: Component 4 — All {total_additional} additional paragraphs preserved (0.15 pts)")
                    total_score += comp4_score
                elif coverage > 0:
                    print(f"PARTIAL: Component 4 — {found_additional}/{total_additional} additional paragraphs found ({comp4_score:.2f} pts)")
                    total_score += comp4_score
                else:
                    print(f"FAIL: Component 4 — No additional paragraphs found in docx")
            elif total_additional == 0:
                # No multi-paragraph slides — internal structure trivially preserved
                print("INFO: Component 4 — No multi-paragraph slides; structure trivially preserved (0.15 pts)")
                total_score += 0.15
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(round(total_score, 4), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


verify_task()
