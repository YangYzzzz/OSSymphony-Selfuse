"""
Initial Setup: Company introduction presentation with 6 slides.
Slide 1 has company name, tagline, and logo - NO animations.
Task ID: impress_ma_086
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_086'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
LOGO_PATH = f'{WORKDIR}/acme_logo.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_logo():
    """Create a simple logo image for the presentation."""
    from PIL import Image, ImageDraw, ImageFont
    img = Image.new('RGBA', (400, 400), (255, 255, 255, 0))
    draw = ImageDraw.Draw(img)
    # Draw a blue circle
    draw.ellipse([50, 50, 350, 350], fill=(0, 102, 204, 255))
    # Draw a white "A" in the center
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 160)
    except (IOError, OSError):
        font = ImageFont.load_default()
    draw.text((140, 80), "A", fill=(255, 255, 255, 255), font=font)
    img.save(LOGO_PATH)
    print(f'Logo created: {LOGO_PATH}')


def create_initial():
    create_logo()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===================== Slide 1: Title / Cover =====================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Company name text box at top
    txBox1 = slide1.shapes.add_textbox(Inches(2), Inches(0.8), Inches(9.333), Inches(1.5))
    tf1 = txBox1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "Acme Corp"
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.runs[0]
    run1.font.name = "Arial"
    run1.font.size = Pt(54)
    run1.font.bold = True
    run1.font.color.rgb = RGBColor(0x00, 0x33, 0x66)
    txBox1.name = "CompanyName"

    # Tagline text box in the middle
    txBox2 = slide1.shapes.add_textbox(Inches(2.5), Inches(2.8), Inches(8.333), Inches(1.2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Innovation Redefined"
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.name = "Arial"
    run2.font.size = Pt(32)
    run2.font.italic = True
    run2.font.color.rgb = RGBColor(0x33, 0x66, 0x99)
    txBox2.name = "Tagline"

    # Logo image at the bottom
    logo_w = Inches(2.5)
    logo_h = Inches(2.5)
    logo_left = (prs.slide_width - logo_w) // 2
    pic = slide1.shapes.add_picture(LOGO_PATH, logo_left, Inches(4.3), logo_w, logo_h)
    pic.name = "Logo"

    # Set slide background
    fill1 = slide1.background.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)

    # ===================== Slide 2: About Us =====================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    title2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
    tf = title2.text_frame
    p = tf.paragraphs[0]
    p.text = "About Acme Corp"
    p.alignment = PP_ALIGN.LEFT
    r = p.runs[0]
    r.font.name = "Arial"
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x00, 0x33, 0x66)

    body2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5))
    tf = body2.text_frame
    tf.word_wrap = True
    texts = [
        "Founded in 2008, Acme Corp has grown from a small startup in San Francisco to a global leader in enterprise technology solutions.",
        "We serve over 2,500 clients across 40 countries, ranging from Fortune 500 companies to innovative startups.",
        "Our mission: empower organizations to achieve operational excellence through cutting-edge software and data-driven insights.",
        "With a team of 3,200+ employees worldwide, we deliver solutions in cloud infrastructure, AI analytics, and cybersecurity."
    ]
    for i, txt in enumerate(texts):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.space_after = Pt(14)
        r = p.runs[0]
        r.font.name = "Arial"
        r.font.size = Pt(18)
        r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ===================== Slide 3: Key Metrics =====================
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    fill3 = slide3.background.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(0x00, 0x33, 0x66)

    title3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
    tf = title3.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Metrics — FY 2025"
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.name = "Arial"
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    metrics = [
        ("$487M", "Annual Revenue"),
        ("2,500+", "Global Clients"),
        ("98.7%", "Client Retention"),
        ("3,200+", "Team Members"),
    ]
    for i, (value, label) in enumerate(metrics):
        left = Inches(0.8 + i * 3.1)
        box = slide3.shapes.add_textbox(left, Inches(2.5), Inches(2.8), Inches(2.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = value
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.name = "Arial"
        r.font.size = Pt(40)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0x4D, 0xC9, 0xF6)
        p2 = tf.add_paragraph()
        p2.text = label
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.runs[0]
        r2.font.name = "Arial"
        r2.font.size = Pt(18)
        r2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    # ===================== Slide 4: Our Services =====================
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    title4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
    tf = title4.text_frame
    p = tf.paragraphs[0]
    p.text = "Our Services"
    p.alignment = PP_ALIGN.LEFT
    r = p.runs[0]
    r.font.name = "Arial"
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x00, 0x33, 0x66)

    services = [
        ("Cloud Infrastructure", "Scalable, secure cloud deployments on AWS, Azure, and GCP with 99.99% uptime SLA."),
        ("AI & Analytics", "Machine learning pipelines and real-time dashboards that turn raw data into actionable insights."),
        ("Cybersecurity", "End-to-end threat detection, compliance automation, and incident response services."),
        ("Digital Transformation", "Strategic consulting to modernize legacy systems and streamline business processes."),
    ]
    for i, (svc, desc) in enumerate(services):
        top = Inches(1.8 + i * 1.3)
        box = slide4.shapes.add_textbox(Inches(0.8), top, Inches(11.7), Inches(1.1))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = svc
        r = p.runs[0]
        r.font.name = "Arial"
        r.font.size = Pt(22)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)
        p2 = tf.add_paragraph()
        p2.text = desc
        r2 = p2.runs[0]
        r2.font.name = "Arial"
        r2.font.size = Pt(16)
        r2.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # ===================== Slide 5: Leadership Team =====================
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    title5 = slide5.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
    tf = title5.text_frame
    p = tf.paragraphs[0]
    p.text = "Leadership Team"
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.name = "Arial"
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x00, 0x33, 0x66)

    leaders = [
        ("Elena Rodriguez", "CEO & Co-Founder", "20+ years in enterprise tech"),
        ("David Park", "CTO", "Former VP Engineering at Cloudflare"),
        ("Priya Sharma", "CFO", "Led IPO at two SaaS companies"),
        ("James Okonkwo", "VP of Sales", "Grew ARR from $50M to $487M"),
    ]
    for i, (name, role, bio) in enumerate(leaders):
        left = Inches(0.5 + i * 3.2)
        box = slide5.shapes.add_textbox(left, Inches(2.0), Inches(2.8), Inches(3.0))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.name = "Arial"
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0x00, 0x33, 0x66)
        p2 = tf.add_paragraph()
        p2.text = role
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.runs[0]
        r2.font.name = "Arial"
        r2.font.size = Pt(16)
        r2.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)
        p3 = tf.add_paragraph()
        p3.text = bio
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.runs[0]
        r3.font.name = "Arial"
        r3.font.size = Pt(14)
        r3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # ===================== Slide 6: Contact Us =====================
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    fill6 = slide6.background.fill
    fill6.solid()
    fill6.fore_color.rgb = RGBColor(0x00, 0x33, 0x66)

    title6 = slide6.shapes.add_textbox(Inches(2), Inches(1.5), Inches(9.333), Inches(1.5))
    tf = title6.text_frame
    p = tf.paragraphs[0]
    p.text = "Get In Touch"
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.name = "Arial"
    r.font.size = Pt(44)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    contact_info = [
        "Email: partnerships@acmecorp.com",
        "Phone: +1 (415) 555-0192",
        "Web: www.acmecorp.com",
        "HQ: 350 Mission Street, San Francisco, CA 94105",
    ]
    contact_box = slide6.shapes.add_textbox(Inches(2.5), Inches(3.2), Inches(8.333), Inches(3.0))
    tf = contact_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(contact_info):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(10)
        r = p.runs[0]
        r.font.name = "Arial"
        r.font.size = Pt(20)
        r.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
