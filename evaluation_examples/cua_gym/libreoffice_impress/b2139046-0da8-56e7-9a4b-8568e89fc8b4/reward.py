"""
Reward Script: Add diagonal watermark 'INTERNAL DRAFT' to master slide
Task ID: impress_gf3_019
Domain: libreoffice_impress
Scoring:
  Component 1: Watermark text 'INTERNAL DRAFT' exists on slide master (0.25)
  Component 2: Font size is 24pt (0.20)
  Component 3: Font color is #AAAAAA (0.20)
  Component 4: Rotation is 45 degrees (0.20)
  Component 5: Watermark is approximately centered on slide (0.15)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_019'


def persist_app_state(domain: str):
    """Save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify that the master slide contains a watermark text box with:
    - Text 'INTERNAL DRAFT'
    - Font size 24pt
    - Font color #AAAAAA
    - 45-degree rotation
    - Approximately centered on slide
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Find the watermark shape on the slide master
    watermark_shape = None
    master = prs.slide_masters[0]
    for shape in master.shapes:
        if hasattr(shape, 'text') and 'INTERNAL DRAFT' in (shape.text or ''):
            watermark_shape = shape
            break

    # Component 1: Watermark text 'INTERNAL DRAFT' exists on slide master (0.25 points)
    try:
        if watermark_shape is not None:
            # Verify the text is exactly 'INTERNAL DRAFT'
            full_text = watermark_shape.text.strip()
            if full_text == 'INTERNAL DRAFT':
                print(f"PASS: Component 1 — Watermark text 'INTERNAL DRAFT' found on master slide (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 — Found text '{full_text}' but expected exactly 'INTERNAL DRAFT'")
        else:
            print("FAIL: Component 1 — No shape with 'INTERNAL DRAFT' text found on slide master")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    if watermark_shape is None:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Font size is 24pt (0.20 points)
    try:
        tf = watermark_shape.text_frame
        runs = [r for p in tf.paragraphs for r in p.runs if (r.text or '').strip()]
        if runs:
            font_size = runs[0].font.size
            expected_size = Pt(24)  # 304800 EMU
            if font_size is not None and font_size == expected_size:
                print(f"PASS: Component 2 — Font size is 24pt ({font_size} EMU) (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 2 — Font size is {font_size} EMU, expected {expected_size} EMU (24pt)")
        else:
            print("FAIL: Component 2 — No text runs found in watermark shape")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Font color is #AAAAAA (0.20 points)
    try:
        runs = [r for p in tf.paragraphs for r in p.runs if (r.text or '').strip()]
        if runs:
            run = runs[0]
            try:
                color_rgb = run.font.color.rgb
                if str(color_rgb).upper() == 'AAAAAA':
                    print(f"PASS: Component 3 — Font color is #AAAAAA (0.20 pts)")
                    total_score += 0.20
                else:
                    print(f"FAIL: Component 3 — Font color is #{color_rgb}, expected #AAAAAA")
            except Exception:
                print("FAIL: Component 3 — Font color is not RGB type or not set")
        else:
            print("FAIL: Component 3 — No text runs found")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Rotation is 45 degrees (0.20 points)
    try:
        rotation = watermark_shape.rotation
        if rotation is not None and abs(rotation - 45.0) < 1.0:
            print(f"PASS: Component 4 — Rotation is {rotation} degrees (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 4 — Rotation is {rotation} degrees, expected 45 degrees")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Watermark is approximately centered on slide (0.15 points)
    try:
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        slide_cx = slide_width / 2.0
        slide_cy = slide_height / 2.0

        shape_cx = watermark_shape.left + watermark_shape.width / 2.0
        shape_cy = watermark_shape.top + watermark_shape.height / 2.0

        # Allow 15% tolerance for centering
        x_tolerance = slide_width * 0.15
        y_tolerance = slide_height * 0.15

        x_ok = abs(shape_cx - slide_cx) <= x_tolerance
        y_ok = abs(shape_cy - slide_cy) <= y_tolerance

        if x_ok and y_ok:
            print(f"PASS: Component 5 — Watermark centered at ({shape_cx:.0f}, {shape_cy:.0f}), "
                  f"slide center ({slide_cx:.0f}, {slide_cy:.0f}) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 5 — Watermark center ({shape_cx:.0f}, {shape_cy:.0f}) "
                  f"not near slide center ({slide_cx:.0f}, {slide_cy:.0f})")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
