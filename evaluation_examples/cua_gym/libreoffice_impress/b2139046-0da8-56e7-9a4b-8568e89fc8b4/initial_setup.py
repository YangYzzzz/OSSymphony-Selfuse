"""
Initial Setup: Create a 20-slide policy draft presentation with no watermark.
Task ID: impress_gf3_019
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_019'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Policy document slide content
    slide_content = [
        {
            "layout": 0,  # Title Slide
            "title": "Corporate Travel & Expense Policy",
            "subtitle": "Effective Date: January 1, 2025\nPrepared by: Finance Department\nVersion 3.2"
        },
        {
            "layout": 1,  # Title + Content
            "title": "Table of Contents",
            "body": "1. Purpose and Scope\n2. Travel Authorization\n3. Booking Guidelines\n4. Expense Categories\n5. Reimbursement Process\n6. Per Diem Rates\n7. International Travel\n8. Receipt Requirements\n9. Corporate Card Usage\n10. Policy Violations"
        },
        {
            "layout": 1,
            "title": "1. Purpose and Scope",
            "body": "This policy establishes guidelines for all business-related travel and expense reimbursement.\n\nApplies to:\n- Full-time employees\n- Part-time employees traveling on company business\n- Contractors with prior VP-level approval\n\nEffective across all regional offices including New York, London, Singapore, and Sydney."
        },
        {
            "layout": 1,
            "title": "2. Travel Authorization Requirements",
            "body": "All domestic travel exceeding $500 requires manager approval.\nInternational travel requires VP-level authorization.\n\nSubmit requests via TravelConnect portal minimum 14 business days in advance.\n\nEmergency travel exceptions must be documented within 48 hours post-travel."
        },
        {
            "layout": 1,
            "title": "3. Booking Guidelines",
            "body": "Use approved booking platform: CorporateTravel Pro\nBook economy class for flights under 6 hours\nBusiness class permitted for flights over 6 hours with director approval\n\nHotel rates must not exceed:\n- Tier 1 cities (NYC, SF, London): $275/night\n- Tier 2 cities: $200/night\n- Tier 3 cities: $150/night"
        },
        {
            "layout": 1,
            "title": "4. Expense Categories",
            "body": "Transportation: Airfare, rail, rental car, rideshare, parking\nLodging: Hotels, serviced apartments (stays > 5 nights)\nMeals: Per diem or actual receipted expenses\nIncidentals: Laundry, tips, phone charges\nClient Entertainment: Pre-approved client meals and events\nConference Fees: Registration, materials, required memberships"
        },
        {
            "layout": 1,
            "title": "5. Reimbursement Process",
            "body": "Submit expense reports within 30 calendar days of travel completion.\nUse ExpenseHub system with digital receipt uploads.\n\nProcessing timeline:\n- Domestic: 10 business days\n- International: 15 business days\n\nDirect deposit to payroll account on file.\nLate submissions (>60 days) require CFO exception approval."
        },
        {
            "layout": 1,
            "title": "6. Per Diem Rates - Domestic",
            "body": "Breakfast: $18\nLunch: $22\nDinner: $38\nIncidentals: $12\nTotal Daily: $90\n\nHigh-cost cities (NYC, SF, DC, Chicago, Boston):\nBreakfast: $24\nLunch: $28\nDinner: $52\nIncidentals: $16\nTotal Daily: $120"
        },
        {
            "layout": 1,
            "title": "7. Per Diem Rates - International",
            "body": "London: GBP 95/day\nTokyo: JPY 15,000/day\nSydney: AUD 140/day\nSingapore: SGD 130/day\nFrankfurt: EUR 85/day\nSao Paulo: BRL 350/day\n\nRates reviewed quarterly against government published rates."
        },
        {
            "layout": 1,
            "title": "8. International Travel",
            "body": "Visa and passport expenses are reimbursable.\nTravel insurance is mandatory — use CorporateShield plan.\n\nCurrency exchange: Use corporate card to minimize fees.\nPersonal exchange losses are not reimbursable.\n\nHealth requirements: Vaccinations reimbursed with medical receipt.\nSecurity briefing required for high-risk destinations."
        },
        {
            "layout": 1,
            "title": "9. Receipt Requirements",
            "body": "Itemized receipts required for all expenses over $25.\nCredit card statements alone are NOT acceptable.\n\nAcceptable formats:\n- Digital photos (minimum 300 DPI)\n- PDF scans\n- Email confirmations with amount and vendor\n\nLost receipt affidavit: Manager sign-off required for amounts $25-$100.\nExpenses over $100 without receipt will not be reimbursed."
        },
        {
            "layout": 1,
            "title": "10. Corporate Card Usage",
            "body": "Platinum Corporate Visa issued to frequent travelers (4+ trips/year).\nPersonal charges strictly prohibited.\n\nMonthly reconciliation deadline: 5th of following month.\nUnreconciled charges after 60 days deducted from payroll.\n\nCard limits:\n- Individual contributor: $5,000\n- Manager: $10,000\n- Director+: $25,000"
        },
        {
            "layout": 1,
            "title": "11. Ground Transportation",
            "body": "Rental cars: Compact/mid-size unless justified by group size.\nInsurance: Decline rental insurance (covered by corporate policy).\n\nRideshare: Standard service only; premium tiers for client travel.\nPersonal vehicle: $0.67/mile reimbursement (IRS 2025 rate).\n\nParking: Airport long-term lot preferred over short-term.\nValet parking only when no self-park option available."
        },
        {
            "layout": 1,
            "title": "12. Client Entertainment",
            "body": "Pre-approval required for expenses over $150.\nDocument attendees, business purpose, and topics discussed.\n\nGuidelines:\n- Alcohol: Limited to 2 drinks per person on company tab\n- Venue: Business-appropriate establishments only\n- Tipping: 18-20% standard; no reimbursement above 25%\n\nAnnual client entertainment budget tracked by cost center."
        },
        {
            "layout": 1,
            "title": "13. Conference & Training Travel",
            "body": "Registration fees: Pre-approved through Learning & Development portal.\nEarly bird rates required when available.\n\nShared accommodations encouraged for team attendance.\nConference materials and required publications reimbursable.\n\nPost-conference knowledge sharing presentation required within 2 weeks."
        },
        {
            "layout": 1,
            "title": "14. Travel Safety & Security",
            "body": "Register all international trips with Global Security team.\nDownload SafeTravel app for real-time alerts.\n\nEmergency contacts:\n- Global Security Hotline: +1-888-555-0199 (24/7)\n- Regional coordinators listed on intranet\n\nEvacuation insurance: Included in CorporateShield plan.\nDo not travel to Level 4 restricted countries without CEO approval."
        },
        {
            "layout": 1,
            "title": "15. Sustainability Guidelines",
            "body": "Prefer rail over air for trips under 300 miles.\nVirtual meetings should be considered before booking travel.\n\nCarbon offset program: Automatic for all booked flights.\nGreen-certified hotels preferred when rate-competitive.\n\nQ4 2025 target: 15% reduction in travel carbon footprint vs 2024."
        },
        {
            "layout": 1,
            "title": "16. Policy Violations",
            "body": "First offense: Written warning and mandatory policy review.\nSecond offense: Travel privileges suspended for 90 days.\nThird offense: Disciplinary action up to termination.\n\nFraudulent expense claims: Immediate termination and legal action.\n\nAll violations reported to Compliance and Internal Audit."
        },
        {
            "layout": 1,
            "title": "17. Frequently Asked Questions",
            "body": "Q: Can I keep frequent flyer miles?\nA: Yes, personal retention permitted for company-booked travel.\n\nQ: What if my flight is cancelled?\nA: Contact CorporateTravel Pro support at +1-800-555-0177.\n\nQ: Can I extend a business trip for personal days?\nA: Yes, with manager approval. Personal days are at own expense."
        },
        {
            "layout": 1,
            "title": "18. Contact & Resources",
            "body": "Travel Department: travel@corporation.com\nExpense Support: expenses@corporation.com\nCorporate Card Services: +1-800-555-0188\n\nPolicy Owner: Sarah Mitchell, VP Finance\nLast Reviewed: December 15, 2024\nNext Review: June 15, 2025\n\nFull policy document available on SharePoint under Finance > Policies"
        },
    ]

    for i, content in enumerate(slide_content):
        layout_idx = content["layout"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        if "title" in content and slide.shapes.title:
            slide.shapes.title.text = content["title"]
            # Style title
            for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                run.font.size = Pt(28) if layout_idx == 0 else Pt(24)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

        if "subtitle" in content:
            slide.placeholders[1].text = content["subtitle"]

        if "body" in content:
            # Find the body placeholder (index 1)
            if 1 in slide.placeholders:
                ph = slide.placeholders[1]
                ph.text = ""
                tf = ph.text_frame
                lines = content["body"].split("\n")
                for j, line in enumerate(lines):
                    if j == 0:
                        tf.paragraphs[0].text = line
                    else:
                        p = tf.add_paragraph()
                        p.text = line
                    # Style the paragraph
                    para = tf.paragraphs[j]
                    for run in para.runs:
                        run.font.size = Pt(14)
                        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
