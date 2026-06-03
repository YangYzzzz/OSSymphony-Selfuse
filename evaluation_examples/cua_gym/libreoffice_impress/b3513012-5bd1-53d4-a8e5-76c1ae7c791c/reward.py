"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 10, LibreOffice Impress labels the image as “Picture 1.” It’s tilting the wrong way—could you show me how to spin that exact picture 15° clockwise so it sits straight?
Generated: 2025-09-10 12:52:23
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
from pptx import Presentation


def normalize_angle(angle):
    """Normalize any angle to the range [-180, 180]."""
    return ((angle + 180) % 360) - 180


def verify_picture_rotation(
    file_path: str,
    slide_index: int = 9,  # zero-based, slide 10 = index 9
    picture_name: str = "Picture 1",
    desired_angle: float = 0.0,  # final correct orientation (straight)
    tolerance: float = 1.0,  # +/- degrees allowed
):
    """Verify that a specific picture on a specific slide has been rotated
    to the desired angle (within tolerance). Progressive scoring is used:
        • 0.4 points for locating the required picture on slide 10.
        • 0.6 points for confirming its rotation is correct within tolerance.
    Returns a float score between 0.0 and 1.0 and prints detailed feedback.
    """

    total_score = 0.0
    max_score = 1.0

    # 1. Load the presentation file (no points for just loading)
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Could not load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # 2. Ensure slide 10 exists
    if len(prs.slides) <= slide_index:
        print(
            f"✗ Slide 10 not found. Presentation only has {len(prs.slides)} slide(s)."
        )
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[slide_index]

    # 3. Locate the target picture by its name
    target_shape = None
    for shape in slide.shapes:
        if shape.name.strip().lower() == picture_name.lower():
            target_shape = shape
            break

    if target_shape is None:
        print(f"✗ Shape named '{picture_name}' not found on slide 10")
        print("REWARD: 0.0")
        return 0.0

    print(f"✓ Found shape '{picture_name}' on slide 10 (0.4 points)")
    total_score += 0.4  # award for correct identification

    # 4. Verify the rotation angle of the picture
    rotation = getattr(target_shape, "rotation", None)
    if rotation is None:
        print("✗ Target shape does not expose a rotation property (unexpected)")
    else:
        normalized_rotation = normalize_angle(rotation)
        print(
            f"Current rotation: {rotation}° (normalized {normalized_rotation}°)"
        )
        if abs(normalized_rotation - desired_angle) <= tolerance:
            print(
                f"✓ Rotation is within ±{tolerance}° of {desired_angle}° (0.6 points)"
            )
            total_score += 0.6
        else:
            print(
                f"✗ Rotation incorrect. Expected ≈{desired_angle}°, got {normalized_rotation}°"
            )

    final_score = round(min(total_score, max_score), 2)
    print(f"REWARD: {final_score}")
    return final_score


# ------------------ EXECUTION ------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/on_slide_10_libreoffice_impress_labels_the_image_as_picture_1_its_tilting_the_wrong_waycould_you_sho_golden.pptx"
    verify_picture_rotation(FILE_PATH)

