"""
Reward Script: Two-tone master slide background
Task ID: impress_ma_039
Domain: libreoffice_impress
Scoring:
  Component 1 (0.35) — Master background color is #ECF0F1
  Component 2 (0.40) — Rectangle on master covering top third, filled #2C3E50
  Component 3 (0.25) — Rectangle dimensions correct (full width, ~1/3 height)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_039'


def persist_app_state(domain: str):
    """Save any unsaved GUI edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(1.0)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Get slide master
    try:
        master = prs.slide_masters[0]
    except Exception as e:
        print(f"CRITICAL: Cannot access slide master: {e}")
        print("REWARD: 0.0")
        return 0.0

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    one_third_height = slide_height / 3.0

    # Component 1: Master background color is #ECF0F1 (0.35 points)
    # In initial state, background is #FFFFFF. Task requires bottom 2/3 to be #ECF0F1.
    try:
        bg_fill = master.background.fill
        if bg_fill.type is not None and bg_fill.type == 1:  # SOLID fill
            bg_color = str(bg_fill.fore_color.rgb).upper()
            if bg_color == "ECF0F1":
                print(f"PASS: Component 1 — Master background is #ECF0F1 (0.35 pts)")
                total_score += 0.35
            else:
                print(f"FAIL: Component 1 — Master background is #{bg_color}, expected #ECF0F1")
        else:
            print(f"FAIL: Component 1 — Master background fill type is {bg_fill.type}, expected SOLID (1)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Rectangle on master filled with #2C3E50 (0.40 points)
    # In initial state, no such rectangle exists on the master.
    try:
        dark_rect = None
        for shape in master.shapes:
            # Look for non-placeholder shapes (AUTO_SHAPE or freeform) with solid fill #2C3E50
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
                try:
                    sfill = shape.fill
                    if sfill.type == 1:  # SOLID
                        color = str(sfill.fore_color.rgb).upper()
                        if color == "2C3E50":
                            dark_rect = shape
                            break
                except Exception:
                    continue

        if dark_rect is not None:
            print(f"PASS: Component 2 — Found rectangle filled with #2C3E50 on master (0.40 pts)")
            total_score += 0.40
        else:
            print(f"FAIL: Component 2 — No rectangle with #2C3E50 fill found on master slide")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Rectangle dimensions — full width, approximately top 1/3 height (0.25 points)
    # Tolerance: 5% relative for position/size checks
    try:
        if dark_rect is not None:
            rect_left = dark_rect.left
            rect_top = dark_rect.top
            rect_width = dark_rect.width
            rect_height = dark_rect.height

            print(f"  Rectangle: left={rect_left}, top={rect_top}, width={rect_width}, height={rect_height}")
            print(f"  Slide: width={slide_width}, height={slide_height}, 1/3 height={one_third_height:.0f}")

            sub_score = 0.0

            # Check top position is at or near 0
            if rect_top <= slide_height * 0.02:  # within 2% of slide height from top
                sub_score += 0.05
                print(f"  PASS: Top position near 0 (top={rect_top})")
            else:
                print(f"  FAIL: Top position not near 0 (top={rect_top})")

            # Check width is approximately full slide width (within 5%)
            width_ratio = rect_width / slide_width if slide_width > 0 else 0
            if width_ratio >= 0.95:
                sub_score += 0.10
                print(f"  PASS: Width is ~full slide width (ratio={width_ratio:.3f})")
            else:
                print(f"  FAIL: Width ratio={width_ratio:.3f}, expected >= 0.95")

            # Check height is approximately 1/3 of slide height (within 15% tolerance)
            height_ratio = rect_height / one_third_height if one_third_height > 0 else 0
            if 0.70 <= height_ratio <= 1.30:
                sub_score += 0.10
                print(f"  PASS: Height ~1/3 slide height (ratio={height_ratio:.3f})")
            else:
                print(f"  FAIL: Height ratio to 1/3={height_ratio:.3f}, expected 0.70-1.30")

            if sub_score > 0:
                print(f"PASS: Component 3 — Rectangle dimensions partially/fully correct ({sub_score:.2f} pts)")
                total_score += sub_score
            else:
                print(f"FAIL: Component 3 — Rectangle dimensions incorrect")
        else:
            print(f"FAIL: Component 3 — No dark rectangle found, cannot check dimensions")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
