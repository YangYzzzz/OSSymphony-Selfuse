"""
Initial Setup: Create a 9-slide presentation with a solid white master background.
Task ID: impress_ma_039
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_039'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard 10x7.5 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Set master slide background to solid white
    master = prs.slide_masters[0]
    bg_fill = master.background.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Slide content themes for a realistic "Dual Tone" themed presentation
    slide_contents = [
        {
            "layout": 0,  # Title Slide
            "title": "Q3 2025 Strategic Review",
            "subtitle": "Dual Tone Consulting Group\nPresented by Elena Rodriguez, Chief Strategy Officer",
        },
        {
            "layout": 1,  # Title + Content
            "title": "Executive Summary",
            "body": "Revenue grew 18% YoY to $42.7M\nClient retention rate improved to 94.3%\nExpanded into 3 new international markets\nLaunched premium advisory service line",
        },
        {
            "layout": 1,
            "title": "Market Analysis",
            "body": "Total addressable market valued at $2.1B\nCompetitor landscape shifting toward digital-first\nEmerging opportunities in healthcare and fintech\nRegulatory changes creating new consulting demand",
        },
        {
            "layout": 1,
            "title": "Financial Performance",
            "body": "Gross margin: 67.2% (up from 61.8%)\nOperating expenses: $28.4M\nEBITDA: $14.3M\nFree cash flow: $9.7M",
        },
        {
            "layout": 1,
            "title": "Client Portfolio Highlights",
            "body": "Secured 12 new enterprise accounts\nAverage deal size increased to $385K\nTop 10 clients represent 43% of revenue\nNPS score improved to 72 from 65",
        },
        {
            "layout": 1,
            "title": "Team & Talent",
            "body": "Headcount grew to 187 employees\nDiversity hiring up 24% year-over-year\nEmployee satisfaction score: 4.3/5.0\nPromoted 15 associates to senior roles",
        },
        {
            "layout": 1,
            "title": "Technology Investments",
            "body": "Deployed AI-powered analytics platform\nMigrated 100% of infrastructure to cloud\nReduced client onboarding time by 40%\nLaunched internal knowledge management system",
        },
        {
            "layout": 1,
            "title": "Risk Assessment",
            "body": "Economic slowdown may impact discretionary spend\nTalent competition intensifying in key markets\nCybersecurity threats require continued investment\nRegulatory compliance costs projected to rise 12%",
        },
        {
            "layout": 1,
            "title": "2026 Strategic Priorities",
            "body": "Expand APAC presence with Singapore office\nLaunch sustainability consulting practice\nAchieve $55M revenue target\nInvest $3.2M in R&D and innovation",
        },
    ]

    for i, sc in enumerate(slide_contents):
        layout_idx = sc["layout"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        if slide.shapes.title:
            slide.shapes.title.text = sc["title"]

        if layout_idx == 0 and "subtitle" in sc:
            # Title slide: placeholder index 1 is subtitle
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = sc["subtitle"]
        elif "body" in sc:
            # Content slide: placeholder index 1 is body
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = sc["body"]

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
