"""
Reward Script: Insert text box on slide 6 with 'Key Takeaways' heading and numbered items
Task ID: impress_tm_081
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): New text box exists on slide 6 with 'Key Takeaways' text
  Component 2 (0.20): 'Key Takeaways' is 24pt bold
  Component 3 (0.25): Three numbered items present with correct text
  Component 4 (0.15): Numbered items are 16pt
  Component 5 (0.15): Numbered items are regular (not bold)
"""

import os
from pptx import Presentation
from pptx.util import Pt

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_081'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[5]  # slide 6 (0-indexed)

    # Find the text box that contains 'Key Takeaways'
    # The initial state has 2 shapes on slide 6; the task adds a new text box
    target_shape = None
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            full_text = shape.text_frame.text
            if 'Key Takeaways' in full_text:
                target_shape = shape
                break

    # Component 1: New text box exists on slide 6 with 'Key Takeaways' (0.25 points)
    try:
        if target_shape is not None:
            print(f"PASS: Component 1 — Found text box with 'Key Takeaways' on slide 6 (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — No text box containing 'Key Takeaways' found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    if target_shape is None:
        # No point checking further components
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    tf = target_shape.text_frame
    paragraphs = tf.paragraphs

    # Find the paragraph that contains 'Key Takeaways'
    heading_para = None
    for para in paragraphs:
        if 'Key Takeaways' in para.text:
            heading_para = para
            break

    # Component 2: 'Key Takeaways' heading is 24pt bold (0.20 points)
    try:
        if heading_para is not None:
            runs = [r for r in heading_para.runs if (r.text or "").strip()]
            if runs:
                heading_run = runs[0]
                is_bold = heading_run.font.bold is True
                # 24pt = 304800 EMU
                size_correct = heading_run.font.size == Pt(24)
                if is_bold and size_correct:
                    print(f"PASS: Component 2 — 'Key Takeaways' is 24pt bold (0.20 pts)")
                    total_score += 0.20
                else:
                    print(f"FAIL: Component 2 — 'Key Takeaways' bold={heading_run.font.bold}, size={heading_run.font.size} (expected bold=True, size={Pt(24)})")
            else:
                print(f"FAIL: Component 2 — No runs found in 'Key Takeaways' paragraph")
        else:
            print(f"FAIL: Component 2 — 'Key Takeaways' paragraph not found")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Three numbered items with correct text (0.25 points)
    expected_items = [
        '1. Embrace change',
        '2. Focus on customers',
        '3. Measure everything',
    ]
    try:
        # Collect all paragraphs after the heading that have non-empty text
        found_items = []
        after_heading = False
        for para in paragraphs:
            if 'Key Takeaways' in para.text:
                after_heading = True
                continue
            if after_heading and para.text.strip():
                found_items.append(para.text.strip())

        matched = 0
        for expected in expected_items:
            for found in found_items:
                if expected.lower() in found.lower():
                    matched += 1
                    break

        if matched == 3:
            print(f"PASS: Component 3 — All 3 numbered items found: {found_items} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 — Found {matched}/3 items. Items: {found_items}, expected: {expected_items}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Numbered items are 16pt (0.15 points)
    try:
        item_paras = []
        after_heading = False
        for para in paragraphs:
            if 'Key Takeaways' in para.text:
                after_heading = True
                continue
            if after_heading and para.text.strip():
                item_paras.append(para)

        size_correct_count = 0
        for para in item_paras:
            runs = [r for r in para.runs if (r.text or "").strip()]
            if runs:
                # 16pt = 203200 EMU
                if runs[0].font.size == Pt(16):
                    size_correct_count += 1

        if len(item_paras) >= 3 and size_correct_count >= 3:
            print(f"PASS: Component 4 — All numbered items are 16pt (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 — {size_correct_count}/{len(item_paras)} items at 16pt")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Numbered items are regular/not bold (0.15 points)
    try:
        not_bold_count = 0
        for para in item_paras:
            runs = [r for r in para.runs if (r.text or "").strip()]
            if runs:
                # bold should be False or None (not bold)
                if runs[0].font.bold is not True:
                    not_bold_count += 1

        if len(item_paras) >= 3 and not_bold_count >= 3:
            print(f"PASS: Component 5 — All numbered items are regular (not bold) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 5 — {not_bold_count}/{len(item_paras)} items are non-bold")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
