"""
Initial Setup: Insert a text box on slide 6 with Key Takeaways heading and numbered items
Task ID: impress_tm_081
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_081'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_body_text(slide, text_lines):
    """Add a text box with body content to a slide."""
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.size = Pt(18)
            run.font.name = "Arial"
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Business Strategy Seminar 2026"
    slide1.placeholders[1].text = "Driving Growth Through Innovation"

    # --- Slide 2: Introduction ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title_box = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf = title_box.text_frame
    tf.paragraphs[0].text = "Introduction"
    tf.paragraphs[0].runs[0].font.size = Pt(32)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.name = "Arial"
    tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
    add_body_text(slide2, [
        "Welcome to our annual strategy seminar where we explore",
        "the latest developments in the business landscape.",
        "This year, we focus on digital transformation and",
        "customer-centric approaches to sustainable growth.",
    ])

    # --- Slide 3: Market Trends ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf = title_box.text_frame
    tf.paragraphs[0].text = "Market Trends"
    tf.paragraphs[0].runs[0].font.size = Pt(32)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.name = "Arial"
    tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
    add_body_text(slide3, [
        "Global e-commerce grew 14% year-over-year reaching $6.3T.",
        "AI adoption increased across 78% of enterprise organizations.",
        "Remote work policies now permanent at 62% of tech firms.",
        "Sustainability reporting became mandatory in 40+ countries.",
    ])

    # --- Slide 4: Competitive Analysis ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide4.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf = title_box.text_frame
    tf.paragraphs[0].text = "Competitive Analysis"
    tf.paragraphs[0].runs[0].font.size = Pt(32)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.name = "Arial"
    tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
    add_body_text(slide4, [
        "Competitor A expanded into Southeast Asian markets with a",
        "30% price reduction strategy targeting mid-market segments.",
        "Competitor B invested $2.1B in AI research and development,",
        "launching three new product lines in Q3 2025.",
    ])

    # --- Slide 5: Growth Strategy ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf = title_box.text_frame
    tf.paragraphs[0].text = "Growth Strategy"
    tf.paragraphs[0].runs[0].font.size = Pt(32)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.name = "Arial"
    tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
    add_body_text(slide5, [
        "Phase 1: Consolidate existing market position by Q2 2026.",
        "Phase 2: Launch premium tier targeting enterprise clients.",
        "Phase 3: International expansion starting with EMEA region.",
        "Phase 4: Strategic partnerships with complementary platforms.",
    ])

    # --- Slide 6: Summary (TITLE ONLY - no text box with takeaways) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide6.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf = title_box.text_frame
    tf.paragraphs[0].text = "Summary"
    tf.paragraphs[0].runs[0].font.size = Pt(32)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.name = "Arial"
    tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
    # NO additional text box - this is the task the agent must perform

    # --- Slide 7: Next Steps ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide7.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf = title_box.text_frame
    tf.paragraphs[0].text = "Next Steps"
    tf.paragraphs[0].runs[0].font.size = Pt(32)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.name = "Arial"
    tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
    add_body_text(slide7, [
        "Schedule follow-up workshops with department leads.",
        "Finalize budget allocation for Q3 initiatives.",
        "Assign project owners for each strategic pillar.",
        "Monthly progress reviews starting June 2026.",
    ])

    # --- Slide 8: Thank You ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide8.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(2))
    tf = title_box.text_frame
    tf.paragraphs[0].text = "Thank You"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].runs[0].font.size = Pt(44)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.name = "Arial"
    tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
    p2 = tf.add_paragraph()
    p2.text = "Questions & Discussion"
    p2.alignment = PP_ALIGN.CENTER
    for run in p2.runs:
        run.font.size = Pt(24)
        run.font.name = "Arial"
        run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
