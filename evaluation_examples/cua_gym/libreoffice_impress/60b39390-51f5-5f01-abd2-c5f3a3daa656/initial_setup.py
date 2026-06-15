"""
Initial Setup: Create a 6-slide Portfolio_Review presentation with empty slide 4
Task ID: impress_stu_083
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_083'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_body_slide(prs, title_text, body_lines):
    """Add a slide with layout 1 (Title + Content), set title and bullet body."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Portfolio Review"
    slide1.placeholders[1].text = "Academic Year 2025-2026"

    # --- Slide 2: About Me ---
    add_title_body_slide(prs, "About Me", [
        "Third-year Computer Science student at Stanford University",
        "Focus areas: Machine Learning, Data Visualization, HCI",
        "Dean's List 2023-2025",
        "Active member of ACM Student Chapter",
        "Teaching assistant for CS 109: Introduction to Probability",
    ])

    # --- Slide 3: Key Projects ---
    add_title_body_slide(prs, "Key Projects", [
        "Sentiment Analysis Dashboard — NLP pipeline with real-time visualization",
        "Campus Wayfinding App — React Native mobile app with indoor positioning",
        "Open-Source Contribution — Added accessibility features to Apache Superset",
        "Research Paper — Co-authored paper on graph neural networks (AAAI 2026)",
        "Hackathon Winner — 1st place at TreeHacks 2025 (Smart Campus IoT)",
    ])

    # --- Slide 4: Competency Assessment (EMPTY — no chart) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Competency Assessment"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 5: Goals ---
    add_title_body_slide(prs, "Goals for Next Semester", [
        "Complete senior thesis on reinforcement learning for robotics",
        "Secure a summer internship at a top AI research lab",
        "Improve public speaking skills through Toastmasters",
        "Learn Rust programming language for systems-level projects",
        "Publish at least one more peer-reviewed paper",
    ])

    # --- Slide 6: Contact ---
    add_title_body_slide(prs, "Contact Information", [
        "Email: alex.rivera@stanford.edu",
        "LinkedIn: linkedin.com/in/alexrivera-cs",
        "GitHub: github.com/arivera-ml",
        "Portfolio: alexrivera.dev",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
