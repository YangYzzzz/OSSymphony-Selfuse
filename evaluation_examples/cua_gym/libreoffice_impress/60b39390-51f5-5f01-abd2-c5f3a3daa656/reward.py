"""
Reward Script: Create a radar/spider chart on slide 4 with self-assessment scores
Task ID: impress_stu_083
Domain: libreoffice_impress
Scoring:
  Component 1: Chart exists on slide 4 (0.20)
  Component 2: Chart is radar type (0.15)
  Component 3: Chart title is 'Skills Self-Assessment' (0.15)
  Component 4: Correct 6 category labels (0.20)
  Component 5: Correct data values for all 6 competencies (0.20)
  Component 6: Value axis scale 0-10 (0.10)
"""

import os
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_083'

# Expected ground truth from task description
EXPECTED_CATEGORIES = [
    'Critical Thinking', 'Communication', 'Technical Skills',
    'Teamwork', 'Time Management', 'Research'
]
EXPECTED_VALUES = [8.0, 7.0, 6.0, 9.0, 5.0, 7.0]


def persist_app_state():
    """Attempt to save any unsaved LibreOffice state."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide4 = prs.slides[3]

    # Find chart shape on slide 4
    chart_shape = None
    for shape in slide4.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Chart exists on slide 4 (0.20 points)
    try:
        if chart_shape is not None:
            print(f"PASS: Component 1 -- Chart found on slide 4 (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 -- No chart found on slide 4")
            # No chart means nothing else to check
            print(f"\nScore: {total_score}/1.0")
            print(f"REWARD: {total_score}")
            return total_score
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    chart = chart_shape.chart

    # Component 2: Chart is radar type (0.15 points)
    # Radar types: RADAR (80), RADAR_FILLED (82), RADAR_MARKERS (81)
    try:
        radar_types = {80, 81, 82}  # RADAR, RADAR_MARKERS, RADAR_FILLED
        chart_type_val = chart.chart_type
        if int(chart_type_val) in radar_types:
            print(f"PASS: Component 2 -- Chart is radar type ({chart_type_val}) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 -- Chart type is {chart_type_val}, expected radar (80/81/82)")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Chart title is 'Skills Self-Assessment' (0.15 points)
    try:
        if chart.has_title:
            title_text = chart.chart_title.text_frame.text.strip()
            if title_text == 'Skills Self-Assessment':
                print(f"PASS: Component 3 -- Chart title is 'Skills Self-Assessment' (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 3 -- Chart title is '{title_text}', expected 'Skills Self-Assessment'")
        else:
            print(f"FAIL: Component 3 -- Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Correct 6 category labels (0.20 points)
    try:
        plot = chart.plots[0]
        actual_cats = [str(c) for c in plot.categories]
        if len(actual_cats) == 6:
            # Check each category matches (case-sensitive)
            matches = sum(1 for a, e in zip(actual_cats, EXPECTED_CATEGORIES) if a.strip() == e)
            if matches == 6:
                print(f"PASS: Component 4 -- All 6 categories match (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 4 -- Only {matches}/6 categories match. Actual: {actual_cats}")
        else:
            print(f"FAIL: Component 4 -- Found {len(actual_cats)} categories, expected 6. Actual: {actual_cats}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Correct data values for all 6 competencies (0.20 points)
    try:
        plot = chart.plots[0]
        if len(plot.series) >= 1:
            actual_vals = list(plot.series[0].values)
            if len(actual_vals) == 6:
                matches = sum(1 for a, e in zip(actual_vals, EXPECTED_VALUES)
                              if abs(float(a) - e) < 0.01)
                if matches == 6:
                    print(f"PASS: Component 5 -- All 6 data values match (0.20 pts)")
                    total_score += 0.20
                else:
                    print(f"FAIL: Component 5 -- Only {matches}/6 values match. Actual: {actual_vals}, Expected: {EXPECTED_VALUES}")
            else:
                print(f"FAIL: Component 5 -- Series has {len(actual_vals)} values, expected 6")
        else:
            print(f"FAIL: Component 5 -- No series found in chart")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # Component 6: Value axis scale 0-10 (0.10 points)
    try:
        va = chart.value_axis
        axis_min = va.minimum_scale
        axis_max = va.maximum_scale
        if axis_min is not None and axis_max is not None:
            if abs(axis_min) < 0.01 and abs(axis_max - 10.0) < 0.01:
                print(f"PASS: Component 6 -- Value axis scale is 0-10 (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 6 -- Axis scale is {axis_min}-{axis_max}, expected 0-10")
        else:
            print(f"FAIL: Component 6 -- Value axis scale not explicitly set (min={axis_min}, max={axis_max})")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
