"""
Initial Setup: Remove all bullet symbols from content textbox on slide 3, increase indent level of all items by one level.
Task ID: osworld_impress_bullet_indent_remove_005
Domain: libreoffice_impress

Creates a 5-slide product comparison deck.
Slide 3 has a content textbox with 4 bullet items at top level (indent level 0) WITH bullet symbols.
The agent must: remove bullet symbols AND increase indent level to 1.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_bullet_indent_remove_005'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_bullet_char_to_para(para, char='•'):
    """Add a visible bullet character to a paragraph via XML."""
    pPr = para._p.get_or_add_pPr()
    # Remove any existing buNone or buChar
    for tag in [qn('a:buNone'), qn('a:buChar'), qn('a:buAutoNum')]:
        existing = pPr.find(tag)
        if existing is not None:
            pPr.remove(existing)
    # Add buChar with the bullet symbol
    buChar = etree.SubElement(pPr, qn('a:buChar'))
    buChar.set('char', char)


def set_para_level(para, level):
    """Set paragraph indent level."""
    para._p.get_or_add_pPr().set('lvl', str(level))


def create_initial():
    prs = Presentation()
    # Use standard 16:9 widescreen
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_layouts = prs.slide_layouts
    # Layout 0 = Title Slide, Layout 1 = Title+Content, Layout 5 = Blank, Layout 2 = Title Only-ish

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(slide_layouts[0])
    slide1.shapes.title.text = "NextGen Product Comparison 2025"
    slide1.placeholders[1].text = "Q2 Strategic Review — Technology Division"

    # ---- Slide 2: Overview ----
    slide2 = prs.slides.add_slide(slide_layouts[1])
    slide2.shapes.title.text = "Overview"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "This report compares four flagship products across key performance indicators."
    p2b = tf2.add_paragraph()
    p2b.text = "Data sourced from internal QA metrics, customer feedback, and market benchmarks."
    p2c = tf2.add_paragraph()
    p2c.text = "All figures reflect Q1 2025 production data."

    # ---- Slide 3: Key Comparison Points (THE TASK SLIDE) ----
    # Has 4 bullet items at level 0 with bullet symbols — agent must remove bullets and increase indent
    slide3 = prs.slides.add_slide(slide_layouts[1])
    slide3.shapes.title.text = "Key Comparison Points"
    tf3 = slide3.placeholders[1].text_frame

    bullet_items = [
        "Performance benchmarks exceed industry average by 23%",
        "Power consumption reduced by 15% compared to previous generation",
        "Customer satisfaction rating increased to 4.7 out of 5.0",
        "Total cost of ownership decreased by 18% over 3-year period",
    ]

    for i, item in enumerate(bullet_items):
        if i == 0:
            para = tf3.paragraphs[0]
            para.text = item
        else:
            para = tf3.add_paragraph()
            para.text = item

        # Set indent level to 0 (top level)
        set_para_level(para, 0)
        # Add bullet character explicitly
        add_bullet_char_to_para(para, char='•')

    # ---- Slide 4: Technical Specifications ----
    slide4 = prs.slides.add_slide(slide_layouts[1])
    slide4.shapes.title.text = "Technical Specifications"
    tf4 = slide4.placeholders[1].text_frame
    specs = [
        "Processor: 3.6 GHz octa-core with neural engine",
        "Memory: 16 GB LPDDR5 unified architecture",
        "Storage: 512 GB NVMe with hardware encryption",
        "Display: 14-inch Retina-class IPS, 2560x1600 resolution",
        "Battery: 72 Wh, rated 14 hours typical workload",
    ]
    for i, spec in enumerate(specs):
        if i == 0:
            tf4.paragraphs[0].text = spec
        else:
            tf4.add_paragraph().text = spec

    # ---- Slide 5: Conclusion & Recommendation ----
    slide5 = prs.slides.add_slide(slide_layouts[1])
    slide5.shapes.title.text = "Conclusion & Recommendation"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Based on the comparative analysis, Product Alpha leads in overall value."
    conc2 = tf5.add_paragraph()
    conc2.text = "Recommended procurement: 500 units for Q3 rollout across APAC region."
    conc3 = tf5.add_paragraph()
    conc3.text = "Follow-up review scheduled for September 2025 to evaluate field performance."

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
