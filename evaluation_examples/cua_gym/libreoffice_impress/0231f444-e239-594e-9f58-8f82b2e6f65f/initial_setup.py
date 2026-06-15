"""
Initial Setup: Create a 5-slide presentation with Blank layout on slide 1
Task ID: impress_gf5_001
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_001'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Blank layout (layout index 6 = Blank in default template) ---
    blank_layout = prs.slide_layouts[6]  # Blank
    slide1 = prs.slides.add_slide(blank_layout)
    # No content on slide 1 - the agent must change layout and add title/subtitle

    # --- Slide 2: Revenue Overview (Title + Content layout) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Revenue Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Total Revenue: $4.2M (up 12% YoY)"
    p = body2.add_paragraph()
    p.text = "North America: $2.1M"
    p.level = 1
    p = body2.add_paragraph()
    p.text = "EMEA: $1.3M"
    p.level = 1
    p = body2.add_paragraph()
    p.text = "APAC: $0.8M"
    p.level = 1
    p = body2.add_paragraph()
    p.text = "Gross margin improved to 68% from 64% in Q2"

    # --- Slide 3: Operational Highlights ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Operational Highlights"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Headcount grew to 342 employees (+15 in Q3)"
    p = body3.add_paragraph()
    p.text = "New office opened in Austin, TX"
    p = body3.add_paragraph()
    p.text = "Customer retention rate: 94.7%"
    p = body3.add_paragraph()
    p.text = "NPS score increased from 62 to 71"

    # --- Slide 4: Product Roadmap ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Product Roadmap"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Q4 Priorities:"
    p = body4.add_paragraph()
    p.text = "Launch mobile app v2.0 (target: Nov 15)"
    p.level = 1
    p = body4.add_paragraph()
    p.text = "Enterprise SSO integration"
    p.level = 1
    p = body4.add_paragraph()
    p.text = "Analytics dashboard redesign"
    p.level = 1
    p = body4.add_paragraph()
    p.text = "API v3 beta release"
    p.level = 1

    # --- Slide 5: Q4 Financial Targets ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Q4 Financial Targets"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Revenue target: $4.8M"
    p = body5.add_paragraph()
    p.text = "Operating expenses budget: $3.1M"
    p = body5.add_paragraph()
    p.text = "EBITDA target: $1.7M"
    p = body5.add_paragraph()
    p.text = "Capital expenditure: $450K (new Austin office buildout)"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
