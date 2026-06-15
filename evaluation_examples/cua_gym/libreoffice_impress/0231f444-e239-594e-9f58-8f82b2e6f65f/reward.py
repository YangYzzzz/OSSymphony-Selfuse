"""
Reward Script: Change slide 1 layout from Blank to Title Slide and add title/subtitle
Task ID: impress_gf5_001
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.30): Slide 1 layout is "Title Slide" (not "Blank")
  - Component 2 (0.30): Title placeholder contains "Q3 Business Review"
  - Component 3 (0.25): Subtitle placeholder contains "Presented by Finance Team"
  - Component 4 (0.15): Presentation still has 5 slides with correct titles on slides 2-5
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_001'


def persist_app_state(domain: str):
    """Best-effort save of any open LibreOffice document."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)

    # Component 1: Slide 1 layout is "Title Slide" (not "Blank") — 0.30 points
    # This FAILS on initial (layout is "Blank") and PASSES on golden (layout is "Title Slide")
    try:
        slide1 = slides[0]
        layout_name = slide1.slide_layout.name if slide1.slide_layout else "Unknown"
        # Accept "Title Slide" or similar names that indicate a title layout
        if "title" in layout_name.lower() and "content" not in layout_name.lower():
            print(f"PASS: Component 1 — Slide 1 layout is '{layout_name}' (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 1 — Expected 'Title Slide' layout, found '{layout_name}'")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Title placeholder contains "Q3 Business Review" — 0.30 points
    # This FAILS on initial (no placeholders on blank slide) and PASSES on golden
    try:
        slide1 = slides[0]
        title_text = None
        for shape in slide1.shapes:
            if shape.has_text_frame and hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                ph_idx = shape.placeholder_format.idx
                # Title placeholder: type CENTER_TITLE (3) or TITLE (1), or idx 0
                if ph_idx == 0:
                    title_text = shape.text_frame.text.strip()
                    break

        if title_text and title_text == "Q3 Business Review":
            print(f"PASS: Component 2 — Title is 'Q3 Business Review' (0.30 pts)")
            total_score += 0.30
        elif title_text:
            print(f"FAIL: Component 2 — Title found but is '{title_text}', expected 'Q3 Business Review'")
        else:
            print(f"FAIL: Component 2 — No title placeholder found on slide 1")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Subtitle placeholder contains "Presented by Finance Team" — 0.25 points
    # This FAILS on initial (no placeholders on blank slide) and PASSES on golden
    try:
        slide1 = slides[0]
        subtitle_text = None
        for shape in slide1.shapes:
            if shape.has_text_frame and hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                ph_idx = shape.placeholder_format.idx
                # Subtitle placeholder: idx 1
                if ph_idx == 1:
                    subtitle_text = shape.text_frame.text.strip()
                    break

        if subtitle_text and subtitle_text == "Presented by Finance Team":
            print(f"PASS: Component 3 — Subtitle is 'Presented by Finance Team' (0.25 pts)")
            total_score += 0.25
        elif subtitle_text:
            print(f"FAIL: Component 3 — Subtitle found but is '{subtitle_text}', expected 'Presented by Finance Team'")
        else:
            print(f"FAIL: Component 3 — No subtitle placeholder found on slide 1")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Presentation still has 5 slides and slides 2-5 titles are unchanged — 0.15 points
    # This FAILS on initial because slide 1 layout change is gated (we check that slide 1
    # has title placeholders AND the other slides are intact together as a compound check).
    # Specifically: we require slide count == 5 AND slides 2-5 have their expected titles
    # AND slide 1 has a title placeholder (anchoring to the change).
    try:
        expected_titles = {
            1: "Revenue Overview",
            2: "Operational Highlights",
            3: "Product Roadmap",
            4: "Q4 Financial Targets",
        }

        if len(slides) != 5:
            print(f"FAIL: Component 4 — Expected 5 slides, found {len(slides)}")
        else:
            # Check that slide 1 has placeholders (task change) AND slides 2-5 are intact
            slide1_has_placeholders = any(
                hasattr(s, 'is_placeholder') and s.is_placeholder
                for s in slides[0].shapes
            )
            if not slide1_has_placeholders:
                print(f"FAIL: Component 4 — Slide 1 has no placeholders (layout not changed)")
            else:
                mismatched = []
                for idx, expected in expected_titles.items():
                    slide = slides[idx]
                    slide_title = None
                    for shape in slide.shapes:
                        if shape.has_text_frame and hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                            if shape.placeholder_format.idx == 0:
                                slide_title = shape.text_frame.text.strip()
                                break
                    if slide_title != expected:
                        mismatched.append((idx + 1, slide_title, expected))

                if len(mismatched) == 0:
                    print(f"PASS: Component 4 — 5 slides present, slides 2-5 unchanged (0.15 pts)")
                    total_score += 0.15
                else:
                    for snum, actual, exp in mismatched:
                        print(f"FAIL: Component 4 — Slide {snum} title is '{actual}', expected '{exp}'")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist any unsaved GUI state before verification
persist_app_state("libreoffice_impress")

# Run verification
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
