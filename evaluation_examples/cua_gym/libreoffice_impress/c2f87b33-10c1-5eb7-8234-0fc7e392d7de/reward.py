"""
Reward Script: ESG Sustainability Report Presentation
Task ID: impress_wf_092
Domain: libreoffice_impress
Scoring:
  C1 (0.15): File exists with exactly 12 slides
  C2 (0.10): Slide 1 title is 'ESG Report 2023'
  C3 (0.10): Slide 2 has 3 arc shapes (gauge semi-circles for E, S, G)
  C4 (0.10): Slides 3-5 have data visualization shapes (charts)
  C5 (0.10): Slide 6 has diversity metrics (pie + bar chart shapes)
  C6 (0.10): Slide 8 has community investment table
  C7 (0.10): Slide 9 has org chart with connected rounded rectangles
  C8 (0.10): Slide 10 has compliance table with checkmarks
  C9 (0.10): Slide 11 has 17 SDG colored rectangle squares
  C10 (0.05): Theme colors #2E7D32 and #5D4037 used in the presentation
"""

import os
import zipfile

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_092'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: File has exactly 12 slides (0.15 points)
    try:
        num_slides = len(prs.slides)
        if num_slides == 12:
            print(f"PASS: Component 1 — 12 slides found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — expected 12 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 1 title contains 'ESG Report 2023' (0.10 points)
    try:
        slide1 = prs.slides[0]
        all_text = ""
        for shape in slide1.shapes:
            if shape.has_text_frame:
                all_text += " " + shape.text_frame.text
        if "ESG Report 2023" in all_text:
            print(f"PASS: Component 2 — 'ESG Report 2023' found on slide 1 (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 — 'ESG Report 2023' not found on slide 1, text: {all_text[:100]}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 2 has 3 arc shapes (gauge semi-circles for E, S, G) (0.10 points)
    try:
        slide2 = prs.slides[1]
        arc_count = 0
        for shape in slide2.shapes:
            # Arcs named 'Arc N' are the semi-circle gauges
            if "Arc" in shape.name or "arc" in shape.name:
                arc_count += 1
        if arc_count >= 3:
            print(f"PASS: Component 3 — {arc_count} arc shapes on slide 2 (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 3 — expected >=3 arc shapes on slide 2, found {arc_count}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slides 3-5 have data visualization shapes (0.10 points)
    # Slide 3: carbon emissions trend (ovals + lines), Slide 4: energy mix (ovals), Slide 5: water stewardship (rectangles as bars)
    try:
        slides_with_viz = 0
        for idx in [2, 3, 4]:  # slides 3, 4, 5
            slide = prs.slides[idx]
            shape_count = len(slide.shapes)
            # Data visualization slides should have many shapes (chart elements)
            if shape_count >= 10:
                slides_with_viz += 1
        if slides_with_viz >= 3:
            print(f"PASS: Component 4 — slides 3-5 all have data visualizations ({slides_with_viz}/3) (0.10 pts)")
            total_score += 0.10
        elif slides_with_viz >= 2:
            print(f"PARTIAL: Component 4 — {slides_with_viz}/3 slides have data visualizations (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 4 — only {slides_with_viz}/3 slides have data visualizations")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slide 6 has diversity metrics (pie + bar chart shapes) (0.10 points)
    try:
        slide6 = prs.slides[5]
        has_pie_indicator = False
        has_bar_indicator = False
        all_text_s6 = ""
        for shape in slide6.shapes:
            if shape.has_text_frame:
                all_text_s6 += " " + shape.text_frame.text
            # Ovals/arcs indicate pie chart; multiple rectangles indicate bar chart
            if "Oval" in shape.name or "Arc" in shape.name:
                has_pie_indicator = True
            if "Rectangle" in shape.name:
                has_bar_indicator = True

        has_diversity_text = "diversity" in all_text_s6.lower() or "gender" in all_text_s6.lower()
        if has_diversity_text and has_pie_indicator and has_bar_indicator:
            print(f"PASS: Component 5 — slide 6 has diversity metrics with pie + bar shapes (0.10 pts)")
            total_score += 0.10
        elif has_diversity_text and (has_pie_indicator or has_bar_indicator):
            print(f"PARTIAL: Component 5 — diversity metrics found but missing chart type (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 5 — slide 6 missing diversity metrics. text_ok={has_diversity_text}, pie={has_pie_indicator}, bar={has_bar_indicator}")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Slide 8 has community investment table (0.10 points)
    try:
        slide8 = prs.slides[7]
        table_found = False
        correct_structure = False
        for shape in slide8.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_found = True
                tbl = shape.table
                rows = len(tbl.rows)
                cols = len(tbl.columns)
                # Should have at least 3 columns (Initiative, Investment, Beneficiaries)
                # and multiple rows
                if rows >= 3 and cols >= 2:
                    header_text = " ".join(tbl.cell(0, c).text for c in range(cols)).lower()
                    if "initiative" in header_text or "investment" in header_text:
                        correct_structure = True

        if table_found and correct_structure:
            print(f"PASS: Component 6 — slide 8 has community investment table (0.10 pts)")
            total_score += 0.10
        elif table_found:
            print(f"PARTIAL: Component 6 — table found on slide 8 but headers don't match (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 6 — no table found on slide 8")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Slide 9 has org chart with connected rounded rectangles (0.10 points)
    try:
        slide9 = prs.slides[8]
        rounded_rect_count = 0
        connector_count = 0
        for shape in slide9.shapes:
            if "Rounded Rectangle" in shape.name or "rounded" in shape.name.lower():
                rounded_rect_count += 1
            if shape.shape_type == 9:  # LINE/CONNECTOR
                connector_count += 1

        if rounded_rect_count >= 5 and connector_count >= 4:
            print(f"PASS: Component 7 — slide 9 org chart: {rounded_rect_count} boxes, {connector_count} connectors (0.10 pts)")
            total_score += 0.10
        elif rounded_rect_count >= 3 and connector_count >= 2:
            print(f"PARTIAL: Component 7 — slide 9 partial org chart: {rounded_rect_count} boxes, {connector_count} connectors (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 7 — slide 9 insufficient org chart shapes: {rounded_rect_count} boxes, {connector_count} connectors")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    # Component 8: Slide 10 has compliance table with checkmarks (0.10 points)
    try:
        slide10 = prs.slides[9]
        table_found = False
        has_checkmarks = False
        for shape in slide10.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_found = True
                tbl = shape.table
                checkmark_count = 0
                for r in range(len(tbl.rows)):
                    for c in range(len(tbl.columns)):
                        cell_text = tbl.cell(r, c).text
                        if "\u2714" in cell_text or "\u2713" in cell_text or "Compliant" in cell_text:
                            checkmark_count += 1
                if checkmark_count >= 3:
                    has_checkmarks = True

        if table_found and has_checkmarks:
            print(f"PASS: Component 8 — slide 10 has compliance table with checkmarks (0.10 pts)")
            total_score += 0.10
        elif table_found:
            print(f"PARTIAL: Component 8 — table on slide 10 but few/no checkmarks (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 8 — no compliance table found on slide 10")
    except Exception as e:
        print(f"ERROR: Component 8 — {e}")

    # Component 9: Slide 11 has 17 SDG colored rectangle squares (0.10 points)
    try:
        slide11 = prs.slides[10]
        rect_count = sum(1 for s in slide11.shapes if "Rectangle" in s.name)
        if rect_count >= 17:
            print(f"PASS: Component 9 — slide 11 has {rect_count} SDG rectangles (0.10 pts)")
            total_score += 0.10
        elif rect_count >= 10:
            print(f"PARTIAL: Component 9 — slide 11 has {rect_count}/17 SDG rectangles (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 9 — slide 11 has only {rect_count} rectangles, expected 17")
    except Exception as e:
        print(f"ERROR: Component 9 — {e}")

    # Component 10: Theme colors #2E7D32 (green) and #5D4037 (brown) used (0.05 points)
    try:
        green_found = False
        brown_found = False
        with zipfile.ZipFile(file_path, 'r') as zf:
            for name in zf.namelist():
                if 'slide' in name and name.endswith('.xml'):
                    content = zf.read(name).decode('utf-8', errors='ignore')
                    if '2E7D32' in content:
                        green_found = True
                    if '5D4037' in content:
                        brown_found = True

        if green_found and brown_found:
            print(f"PASS: Component 10 — both theme colors #2E7D32 and #5D4037 found (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 10 — green={green_found}, brown={brown_found}")
    except Exception as e:
        print(f"ERROR: Component 10 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point — run on VM
file_path = f'{WORKDIR}/Desktop/ESG_Report.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
