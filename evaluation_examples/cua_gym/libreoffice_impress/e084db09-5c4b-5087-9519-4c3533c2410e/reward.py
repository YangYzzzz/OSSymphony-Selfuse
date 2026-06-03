"""
FINAL REWARD SCRIPT - SUCCESS
Task: Convert all footnotes to endnotes.
Generated: 2025-10-17 12:05:46
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
import re
from pptx import Presentation


def verify_footnotes_to_endnotes(file_path: str) -> float:
    """Verify that all footnotes in the presentation were converted to endnotes.

    Scoring (progressive, max 1.0):
    1. Detect an "Endnotes" slide anywhere in the deck ........................ 0.3
    2. The Endnotes slide is the LAST slide .................................... 0.1
    3. Endnotes slide contains at least one numbered end-note line (e.g. "1.").. 0.3
    4. No other slide (except the Endnotes slide) contains residual footnote
       lines starting with a number and a dot (heuristic) ...................... 0.3
    """

    max_score = 1.0
    score = 0.0

    # ---------- 1) Load presentation ----------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0  # Nothing to score

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        return 0.0

    total_slides = len(prs.slides)
    print(f"Total slides detected: {total_slides}")

    # ---------- 2) Locate the Endnotes slide ----------
    endnotes_slide_idx = None  # 0-based index
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                if shape.text.strip().lower() == "endnotes":
                    endnotes_slide_idx = idx
                    break
        if endnotes_slide_idx is not None:
            break

    if endnotes_slide_idx is not None:
        print(f"✓ Found 'Endnotes' slide at position {endnotes_slide_idx + 1}")
        score += 0.3
        if endnotes_slide_idx == total_slides - 1:
            print("✓ 'Endnotes' slide is the last slide")
            score += 0.1
        else:
            print("✗ 'Endnotes' slide is not the last slide")
    else:
        print("✗ No slide titled 'Endnotes' found – cannot continue full verification")
        return min(score, max_score)

    # ---------- 3) Verify numbered lines in Endnotes slide ----------
    numbered_lines = []
    end_slide = prs.slides[endnotes_slide_idx]
    for shape in end_slide.shapes:
        if hasattr(shape, "text") and shape.text:
            for line in shape.text.splitlines():
                if re.match(r"\s*\d+\.\s+", line):
                    numbered_lines.append(line.strip())

    if numbered_lines:
        print(f"✓ Detected numbered endnote lines: {numbered_lines}")
        score += 0.3
    else:
        print("✗ No numbered endnote lines found on the 'Endnotes' slide")

    # ---------- 4) Ensure no residual footnotes on other slides ----------
    residuals = []
    for idx, slide in enumerate(prs.slides):
        if idx == endnotes_slide_idx:
            continue  # Skip Endnotes slide itself
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                for line in shape.text.splitlines():
                    if re.match(r"\s*\d+\.\s+", line):  # looks like a footnote marker
                        residuals.append((idx + 1, line.strip()))

    if residuals:
        print("✗ Potential residual footnotes found on other slides:")
        for slide_no, line in residuals:
            print(f"   • Slide {slide_no}: '{line}'")
    else:
        print("✓ No residual footnote lines found on other slides")
        score += 0.3

    final_score = min(score, max_score)
    print(f"Final computed score: {final_score}")
    return final_score


if __name__ == "__main__":
    # Path provided by the task context
    FILE_PATH = "/home/user/convert_all_footnotes_to_endnotes.pptx"

    reward_value = verify_footnotes_to_endnotes(FILE_PATH)
    print(f"REWARD: {reward_value}")

