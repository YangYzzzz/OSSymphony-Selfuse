"""
FINAL REWARD SCRIPT - SUCCESS
Task: I'm trying to make the logo on the first slide into a neat square shape so that it's centered perfectly. Any tips on how I can crop it like that in Impress?
Generated: 2025-08-07 08:54:32
Status: success
Model: o4-mini
Total Steps: 1
"""

import os
from pptx import Presentation
from pptx.shapes.picture import Picture

def verify_logo_cropping_centering(file_path):
    """
    Verifies that the first slide of the presentation at file_path contains a picture,
    that the picture is cropped to a square (width == height), and that it is centered
    both horizontally and vertically on the slide. Returns a float reward between 0.0 and 1.0.
    """
    print("Checking task completion for logo cropping and centering...")
    score = 0.0
    max_score = 1.0

    # Requirement 1: File existence and successful load (0.2 points)
    try:
        if not os.path.exists(file_path):
            print(f"✗ File not found: {file_path}")
            print(f"REWARD: {score}")
            return score
        prs = Presentation(file_path)
        print("✓ File exists and presentation loaded (0.2 points)")
        score += 0.2
    except Exception as e:
        print(f"✗ Error accessing or loading file: {e}")
        print(f"REWARD: {score}")
        return score

    # Requirement 2: Presence of a picture on the first slide (0.2 points)
    try:
        slide = prs.slides[0]
        pics = [shape for shape in slide.shapes if isinstance(shape, Picture)]
        if not pics:
            print("✗ No picture shapes found on the first slide.")
            print(f"REWARD: {score}")
            return score
        pic = pics[0]
        print(f"✓ Found {len(pics)} picture shape(s) on first slide (0.2 points)")
        score += 0.2
    except Exception as e:
        print(f"✗ Error checking picture shapes: {e}")
        print(f"REWARD: {score}")
        return score

    # Requirement 3: Verify picture is square (width == height) (0.3 points)
    width = pic.width
    height = pic.height
    print(f"Picture dimensions: width={width}, height={height}")
    if width == height:
        print("✓ Picture is square (0.3 points)")
        score += 0.3
    else:
        print("✗ Picture is not square. Consider cropping to equal width and height.")

    # Requirement 4: Verify picture is centered on the slide (0.3 points)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    expected_left = (slide_w - width) // 2
    expected_top = (slide_h - height) // 2
    actual_left = pic.left
    actual_top = pic.top
    print(f"Slide size: width={slide_w}, height={slide_h}")
    print(f"Expected position: left={expected_left}, top={expected_top}")
    print(f"Actual position: left={actual_left}, top={actual_top}")
    if actual_left == expected_left and actual_top == expected_top:
        print("✓ Picture is centered (0.3 points)")
        score += 0.3
    else:
        print("✗ Picture is not centered perfectly. Consider aligning center horizontally and vertically.")

    # Finalize and cap score
    final_score = min(score, max_score)
    print(f"Final score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score

# Execute verification when run as a script
if __name__ == '__main__':
    file_path = '/home/user/im_trying_to_make_the_logo_on_the_first_slide_into_a_neat_square_shape_so_that_its_centered_perfectl.pptx'
    verify_logo_cropping_centering(file_path)
