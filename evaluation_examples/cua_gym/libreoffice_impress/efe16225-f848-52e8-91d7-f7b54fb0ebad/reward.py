"""
Reward Script: Insert a 2x3 table on slide 1 with team contact info, positioned in bottom-right quarter.
Task ID: impress_tct_033
Domain: libreoffice_impress
Scoring:
  Component 1 (0.40): Table exists on slide 1 with 2 columns and 3 rows
  Component 2 (0.35): Table positioned in bottom-right quarter of the slide
  Component 3 (0.25): Table contains contact-related content (names/emails/phone)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_033'

def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 1 slide
    if len(prs.slides) < 1:
        print("FAIL: Presentation has no slides")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]
    slide_width = prs.slide_width   # EMU
    slide_height = prs.slide_height  # EMU

    # Find table shapes on slide 1
    table_shapes = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_shapes.append(shape)

    if not table_shapes:
        print("FAIL: No table found on slide 1")
        print("REWARD: 0.0")
        return 0.0

    # Use the first (or only) table found
    table_shape = table_shapes[0]
    table = table_shape.table

    num_rows = len(table.rows)
    num_cols = len(table.columns)

    # Component 1: Table exists on slide 1 with 2 columns and 3 rows (0.40 points)
    try:
        if num_cols == 2 and num_rows == 3:
            print(f"PASS: Component 1 — Table has {num_rows} rows x {num_cols} cols (0.40 pts)")
            total_score += 0.40
        else:
            print(f"FAIL: Component 1 — Expected 3 rows x 2 cols, found {num_rows} rows x {num_cols} cols")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Table positioned in bottom-right quarter (0.35 points)
    # Slide is 10x7.5 inches. Bottom-right quarter means:
    #   left >= slide_width / 2 (>= 5.0 inches)
    #   top >= slide_height / 2 (>= 3.75 inches)
    # We use a tolerance: left >= 4.5in, top >= 3.5in to allow reasonable positioning
    try:
        left_in = table_shape.left / 914400.0
        top_in = table_shape.top / 914400.0
        width_in = table_shape.width / 914400.0
        height_in = table_shape.height / 914400.0

        half_w = slide_width / 2.0 / 914400.0  # 5.0 inches
        half_h = slide_height / 2.0 / 914400.0  # 3.75 inches

        # Table left edge should be at or past the horizontal midpoint (with tolerance)
        # Table top edge should be at or past the vertical midpoint (with tolerance)
        in_right_half = left_in >= (half_w - 0.75)  # >= 4.25 inches
        in_bottom_half = top_in >= (half_h - 0.5)   # >= 3.25 inches

        if in_right_half and in_bottom_half:
            print(f"PASS: Component 2 — Table at left={left_in:.2f}in, top={top_in:.2f}in "
                  f"(bottom-right quarter) (0.35 pts)")
            total_score += 0.35
        else:
            print(f"FAIL: Component 2 — Table at left={left_in:.2f}in, top={top_in:.2f}in. "
                  f"Expected left>={half_w - 0.75:.2f}in, top>={half_h - 0.5:.2f}in "
                  f"for bottom-right quarter")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Table contains contact-related content (0.25 points)
    # Check that at least some cells contain text that looks like contact info
    # (names, emails, phone numbers, or contact-related headers)
    try:
        all_cell_text = []
        for r in range(num_rows):
            for c in range(num_cols):
                cell_text = table.cell(r, c).text.strip()
                all_cell_text.append(cell_text)

        combined = " ".join(all_cell_text).lower()

        # Check for contact-related content indicators
        has_contact_headers = any(kw in combined for kw in
            ["name", "email", "phone", "contact", "role", "title", "department"])
        has_at_sign = "@" in combined  # email indicator
        has_multiple_entries = sum(1 for t in all_cell_text if len(t) > 0) >= 4  # at least 4 non-empty cells

        # Need at least 2 of 3 indicators
        indicators = sum([has_contact_headers, has_at_sign, has_multiple_entries])

        if indicators >= 2:
            print(f"PASS: Component 3 — Table contains contact-related content "
                  f"(headers={has_contact_headers}, emails={has_at_sign}, "
                  f"entries={has_multiple_entries}) (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 — Table content does not appear to be contact info. "
                  f"Indicators: headers={has_contact_headers}, emails={has_at_sign}, "
                  f"entries={has_multiple_entries}. Cell contents: {all_cell_text}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
