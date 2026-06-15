"""
Initial Setup: Create a 3-slide Team_Intro presentation with slide 1 having
a title and subtitle, lower-right area empty for future table insertion.
Task ID: impress_tct_033
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_033'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard slide size: 10 x 7.5 inches
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
    slide1.shapes.title.text = "Meet Our Team"
    slide1.placeholders[1].text = "Engineering & Product Division - Q2 2025"

    # --- Slide 2: Team Members ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide2.shapes.title.text = "Our Team Members"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Sarah Chen - Lead Engineer"
    p2 = body2.add_paragraph()
    p2.text = "Marcus Johnson - Product Manager"
    p2.level = 0
    p3 = body2.add_paragraph()
    p3.text = "Aisha Patel - UX Designer"
    p3.level = 0
    p4 = body2.add_paragraph()
    p4.text = "David Kim - Backend Developer"
    p4.level = 0
    p5 = body2.add_paragraph()
    p5.text = "Elena Rodriguez - QA Lead"
    p5.level = 0

    # --- Slide 3: Mission & Values ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide3.shapes.title.text = "Our Mission"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "We build tools that empower teams to collaborate effectively."
    q1 = body3.add_paragraph()
    q1.text = "Core Values:"
    q1.level = 0
    for run in q1.runs:
        run.font.bold = True
    q2 = body3.add_paragraph()
    q2.text = "Innovation - We push boundaries every day"
    q2.level = 1
    q3 = body3.add_paragraph()
    q3.text = "Collaboration - We succeed together"
    q3.level = 1
    q4 = body3.add_paragraph()
    q4.text = "Excellence - We hold ourselves to high standards"
    q4.level = 1

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
