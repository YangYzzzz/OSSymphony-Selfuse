"""
Reward Script: Traffic light textboxes on slide 2 with correct colors
Task ID: osworld_impress_textbox_colors_multiple_008
Domain: libreoffice_impress
Scoring:
  Component 1: 3 textboxes present on slide 2 (0.30 pts)
  Component 2: Correct text labels and font colors per textbox (0.50 pts, ~0.167 each)
  Component 3: Vertical arrangement with even spacing (0.20 pts)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_textbox_colors_multiple_008'

# Expected traffic light entries: (text, hex_color_uppercase)
EXPECTED_TEXTBOXES = [
    ('On Track', '28A745'),
    ('At Risk', 'FFC107'),
    ('Off Track', 'DC3545'),
]


def get_textbox_shapes(slide):
    """Return all TEXT_BOX (type 17) shapes from a slide, excluding group members."""
    result = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            result.append(shape)
    return result


def get_run_color(run):
    """Return the RGB hex string of a run's font color, or None if not set."""
    try:
        if run.font.color.type is not None:
            return str(run.font.color.rgb).upper()
    except Exception:
        pass
    return None


def get_shape_text(shape):
    """Return the full text of a shape (strips whitespace)."""
    if shape.has_text_frame:
        return shape.text_frame.text.strip()
    return ''


def verify_task(file_path):
    """
    Verify task completion: 3 color-coded textboxes on slide 2 with
    correct text, colors, and even vertical spacing.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Verify we have at least 2 slides (slide 2 must exist)
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has only {len(prs.slides)} slide(s); expected at least 2")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]  # 0-indexed — slide 2

    # ----- Component 1: 3 textboxes present on slide 2 (0.30 pts) -----
    try:
        textboxes = get_textbox_shapes(slide2)
        num_tb = len(textboxes)
        if num_tb == 3:
            print(f"PASS: Component 1 — exactly 3 textboxes found on slide 2 (0.30 pts)")
            total_score += 0.30
        elif num_tb > 0:
            print(f"FAIL: Component 1 — found {num_tb} textbox(es) on slide 2; expected exactly 3")
        else:
            print(f"FAIL: Component 1 — no textboxes found on slide 2")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # ----- Component 2: Correct text + color for each textbox (0.50 pts total) -----
    # Each matched entry awards ~0.1667 pts; we use 0.50/3 per entry
    POINTS_PER_ENTRY = round(0.50 / 3, 4)  # ~0.1667
    try:
        textboxes = get_textbox_shapes(slide2)
        # Build a dict: text -> list of colors found in that textbox
        found_map = {}
        for tb in textboxes:
            tb_text = get_shape_text(tb)
            if tb.has_text_frame:
                colors_in_tb = []
                for para in tb.text_frame.paragraphs:
                    for run in para.runs:
                        if (run.text or '').strip():
                            c = get_run_color(run)
                            if c:
                                colors_in_tb.append(c)
                found_map[tb_text] = colors_in_tb

        matched = 0
        for expected_text, expected_color in EXPECTED_TEXTBOXES:
            if expected_text in found_map:
                colors = found_map[expected_text]
                if expected_color in colors:
                    print(f"PASS: Component 2 — '{expected_text}' found with color #{expected_color} ({POINTS_PER_ENTRY} pts)")
                    total_score += POINTS_PER_ENTRY
                    matched += 1
                else:
                    print(f"FAIL: Component 2 — '{expected_text}' found but color is {colors}; expected #{expected_color}")
            else:
                available_texts = list(found_map.keys())
                print(f"FAIL: Component 2 — '{expected_text}' not found in slide 2 textboxes (found: {available_texts})")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # ----- Component 3: Vertical arrangement with even spacing (0.20 pts) -----
    try:
        textboxes = get_textbox_shapes(slide2)
        if len(textboxes) == 3:
            # Sort by top position to check ordering
            sorted_tbs = sorted(textboxes, key=lambda s: s.top)
            tops = [tb.top for tb in sorted_tbs]

            # Check all tops are strictly increasing (vertical arrangement)
            strictly_increasing = tops[0] < tops[1] < tops[2]

            # Check even spacing: gap1 == gap2 within 5% relative tolerance
            gap1 = tops[1] - tops[0]
            gap2 = tops[2] - tops[1]
            if gap1 > 0 and gap2 > 0:
                spacing_ratio = abs(gap1 - gap2) / max(gap1, gap2)
                even_spacing = spacing_ratio <= 0.05
            else:
                even_spacing = False

            if strictly_increasing and even_spacing:
                print(f"PASS: Component 3 — textboxes vertically arranged with even spacing "
                      f"(tops={tops}, gap1={gap1}, gap2={gap2}) (0.20 pts)")
                total_score += 0.20
            elif strictly_increasing:
                print(f"FAIL: Component 3 — textboxes are ordered top-to-bottom but spacing is uneven "
                      f"(gap1={gap1}, gap2={gap2}, ratio={spacing_ratio:.3f})")
            else:
                print(f"FAIL: Component 3 — textboxes are NOT in strict vertical order (tops={tops})")
        else:
            print(f"FAIL: Component 3 — cannot check spacing, expected 3 textboxes (found {len(textboxes)})")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 4)
    print(f"\nScore: {total_score:.4f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in this env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
