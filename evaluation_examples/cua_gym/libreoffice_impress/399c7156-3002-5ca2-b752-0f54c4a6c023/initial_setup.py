"""
Initial Setup: Project dashboard presentation - Slide 2 'Project Status' is empty (no status textboxes)
Task ID: osworld_impress_textbox_colors_multiple_008
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_textbox_colors_multiple_008'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide layouts: 0=Title Slide, 1=Title+Content, 5=Blank, 6=Title Only
    # Slide 1: Title Slide — "Project Dashboard"
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Project Dashboard"
    slide1.placeholders[1].text = "Q2 2025 Executive Overview"

    # Slide 2: "Project Status" — title only, NO status textboxes (task is to add them)
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide2.shapes.title.text = "Project Status"
    # No textboxes added — agent must add them

    # Slide 3: "Timeline" with content
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Project Timeline"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Phase 1: Discovery & Planning"
    p3b = tf3.add_paragraph()
    p3b.text = "Phase 2: Design & Development"
    p3c = tf3.add_paragraph()
    p3c.text = "Phase 3: Testing & QA"
    p3d = tf3.add_paragraph()
    p3d.text = "Phase 4: Deployment & Handoff"

    # Slide 4: "Budget Overview" with content
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Budget Overview"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Total Budget: $2,450,000"
    p4b = tf4.add_paragraph()
    p4b.text = "Spent to Date: $1,123,750"
    p4c = tf4.add_paragraph()
    p4c.text = "Remaining: $1,326,250"
    p4d = tf4.add_paragraph()
    p4d.text = "Forecast Variance: +$48,000"

    # Slide 5: "Team & Resources" with content
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Team & Resources"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Project Lead: Alexandra Rivera"
    p5b = tf5.add_paragraph()
    p5b.text = "Engineering: 8 FTEs"
    p5c = tf5.add_paragraph()
    p5c.text = "Design: 3 FTEs"
    p5d = tf5.add_paragraph()
    p5d.text = "QA: 4 FTEs"

    # Slide 6: "Next Steps" with content
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Next Steps"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Finalize sprint planning for Q3"
    p6b = tf6.add_paragraph()
    p6b.text = "Conduct stakeholder review meeting"
    p6c = tf6.add_paragraph()
    p6c.text = "Submit compliance documentation"
    p6d = tf6.add_paragraph()
    p6d.text = "Begin user acceptance testing"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open the file in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
