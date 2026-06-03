"""
Initial Setup: Create a Renewable Energy presentation with 8 slides.
Slide 5 has title 'Solar Energy: Pros and Cons' with empty content area.
Task ID: impress_teach_024
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_024'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Helper to add a text box with title text
    def add_title_content_slide(title_text, body_text):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
        slide.shapes.title.text = title_text
        slide.placeholders[1].text = body_text
        return slide

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    slide1.shapes.title.text = "Renewable Energy"
    slide1.placeholders[1].text = "A Comprehensive Overview\nPrepared by the Energy Research Institute\nMarch 2025"

    # --- Slide 2: Introduction ---
    add_title_content_slide(
        "Introduction to Renewable Energy",
        "Renewable energy comes from sources that are naturally replenished on a human timescale.\n\n"
        "Key sources include solar, wind, hydroelectric, geothermal, and biomass energy.\n\n"
        "Global renewable energy capacity reached 3,372 GW in 2024, a 12% increase from the previous year."
    )

    # --- Slide 3: Types of Renewable Energy ---
    add_title_content_slide(
        "Types of Renewable Energy",
        "Solar Energy - Photovoltaic cells and concentrated solar power\n"
        "Wind Energy - Onshore and offshore wind turbines\n"
        "Hydroelectric - Dam-based and run-of-river systems\n"
        "Geothermal - Earth's internal heat for electricity and heating\n"
        "Biomass - Organic materials converted to fuel or electricity"
    )

    # --- Slide 4: Wind Energy Overview ---
    add_title_content_slide(
        "Wind Energy Overview",
        "Global installed wind capacity: 906 GW (2024)\n\n"
        "Onshore wind farms account for 85% of total capacity.\n\n"
        "Leading countries: China (395 GW), United States (148 GW), Germany (69 GW)\n\n"
        "Average cost of onshore wind: $0.033/kWh, competitive with fossil fuels."
    )

    # --- Slide 5: Solar Energy Pros and Cons (EMPTY CONTENT) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title text box at top
    title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Solar Energy: Pros and Cons"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    # No comparison text boxes - this is the empty content area for the task

    # --- Slide 6: Hydroelectric Power ---
    add_title_content_slide(
        "Hydroelectric Power",
        "Largest renewable electricity source globally, providing 16% of world electricity.\n\n"
        "Three Gorges Dam (China): 22.5 GW capacity, world's largest power station.\n\n"
        "Advantages: Low operating costs, long lifespan, flood control.\n\n"
        "Environmental concerns include habitat disruption and displacement of communities."
    )

    # --- Slide 7: Geothermal Energy ---
    add_title_content_slide(
        "Geothermal Energy",
        "Harnesses heat from Earth's interior for electricity generation and direct heating.\n\n"
        "Global capacity: 16.1 GW across 30 countries.\n\n"
        "Iceland generates 25% of its electricity from geothermal sources.\n\n"
        "Enhanced geothermal systems (EGS) could expand viable locations significantly."
    )

    # --- Slide 8: Future Outlook ---
    add_title_content_slide(
        "Future Outlook",
        "Renewables projected to generate 50% of global electricity by 2030.\n\n"
        "Battery storage costs declining 15% annually, solving intermittency challenges.\n\n"
        "Green hydrogen emerging as key energy carrier for hard-to-decarbonize sectors.\n\n"
        "Investment in renewable energy surpassed $500 billion in 2024."
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
