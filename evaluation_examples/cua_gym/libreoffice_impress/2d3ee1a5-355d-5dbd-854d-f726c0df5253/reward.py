"""
Reward Script: Two-column Advantages/Disadvantages layout on slide 5
Task ID: impress_teach_024
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.25): Left text box with 'Advantages' title exists on slide 5
  - Component 2 (0.25): 'Advantages' title is green (#2E7D32)
  - Component 3 (0.25): Right text box with 'Disadvantages' title exists on slide 5
  - Component 4 (0.25): 'Disadvantages' title is red (#C62828)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_024'


def persist_app_state(domain: str):
    """Save any unsaved GUI edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
    except ImportError:
        print("CRITICAL: python-pptx not installed")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, expected at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # 0-indexed, slide 5
    slide_mid_x = prs.slide_width / 2

    # Find all text boxes on slide 5 (exclude placeholders that are part of the layout)
    # We look for text boxes containing 'Advantages' or 'Disadvantages' as their first paragraph text
    advantages_box = None
    disadvantages_box = None

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # Get the first non-empty paragraph text
        first_text = ""
        first_run_color = None
        first_run_bold = None
        for para in shape.text_frame.paragraphs:
            stripped = para.text.strip()
            if stripped:
                first_text = stripped
                # Get color from runs
                for run in para.runs:
                    if (run.text or "").strip():
                        try:
                            if run.font.color.type is not None:
                                first_run_color = str(run.font.color.rgb)
                        except Exception:
                            pass
                        first_run_bold = run.font.bold
                        break
                break

        if first_text.lower() == 'advantages':
            advantages_box = {
                'shape': shape,
                'text': first_text,
                'color': first_run_color,
                'bold': first_run_bold,
                'left': shape.left,
                'center_x': shape.left + shape.width / 2,
            }
        elif first_text.lower() == 'disadvantages':
            disadvantages_box = {
                'shape': shape,
                'text': first_text,
                'color': first_run_color,
                'bold': first_run_bold,
                'left': shape.left,
                'center_x': shape.left + shape.width / 2,
            }

    # Component 1: Left text box with 'Advantages' exists on slide 5 (0.25 points)
    try:
        if advantages_box is not None:
            # Verify it's on the left side of the slide
            if advantages_box['center_x'] < slide_mid_x:
                print(f"PASS: Component 1 - 'Advantages' text box found on left side of slide 5 (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 - 'Advantages' text box exists but is NOT on the left side (center_x={advantages_box['center_x']}, mid={slide_mid_x})")
        else:
            print(f"FAIL: Component 1 - No text box with 'Advantages' text found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: 'Advantages' title is green #2E7D32 (0.25 points)
    try:
        if advantages_box is not None and advantages_box['color'] is not None:
            actual_color = advantages_box['color'].upper()
            if actual_color == '2E7D32':
                print(f"PASS: Component 2 - 'Advantages' color is #2E7D32 (green) (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 2 - 'Advantages' color is #{actual_color}, expected #2E7D32")
        else:
            if advantages_box is None:
                print(f"FAIL: Component 2 - No 'Advantages' text box found, cannot check color")
            else:
                print(f"FAIL: Component 2 - 'Advantages' text has no explicit color set")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Right text box with 'Disadvantages' exists on slide 5 (0.25 points)
    try:
        if disadvantages_box is not None:
            # Verify it's on the right side of the slide
            if disadvantages_box['center_x'] > slide_mid_x:
                print(f"PASS: Component 3 - 'Disadvantages' text box found on right side of slide 5 (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 3 - 'Disadvantages' text box exists but is NOT on the right side (center_x={disadvantages_box['center_x']}, mid={slide_mid_x})")
        else:
            print(f"FAIL: Component 3 - No text box with 'Disadvantages' text found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: 'Disadvantages' title is red #C62828 (0.25 points)
    try:
        if disadvantages_box is not None and disadvantages_box['color'] is not None:
            actual_color = disadvantages_box['color'].upper()
            if actual_color == 'C62828':
                print(f"PASS: Component 4 - 'Disadvantages' color is #C62828 (red) (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 4 - 'Disadvantages' color is #{actual_color}, expected #C62828")
        else:
            if disadvantages_box is None:
                print(f"FAIL: Component 4 - No 'Disadvantages' text box found, cannot check color")
            else:
                print(f"FAIL: Component 4 - 'Disadvantages' text has no explicit color set")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
