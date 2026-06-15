"""
Initial Setup: Create a 10-slide workshop presentation with no transitions or slide numbers.
Task ID: impress_gf4_009
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_009'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                alignment=PP_ALIGN.LEFT, color=None):
    """Helper to add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tf


def add_bullet_text(text_frame, text, level=0, font_size=16):
    """Add a bullet point paragraph."""
    p = text_frame.add_paragraph()
    p.level = level
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    return p


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ========== Slide 1: Title Slide ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    add_textbox(slide1, Inches(1), Inches(1.5), Inches(11), Inches(2),
                "Data Analytics Workshop 2025", font_size=40, bold=True,
                alignment=PP_ALIGN.CENTER, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_textbox(slide1, Inches(1), Inches(4), Inches(11), Inches(1),
                "Building Insights from Raw Data", font_size=24,
                alignment=PP_ALIGN.CENTER, color=RGBColor(0xCC, 0xDD, 0xEE))
    add_textbox(slide1, Inches(1), Inches(5.5), Inches(11), Inches(1),
                "Presented by Elena Vasquez  |  March 18, 2025  |  TechHub Conference Center",
                font_size=16, alignment=PP_ALIGN.CENTER,
                color=RGBColor(0x99, 0xBB, 0xDD))

    # ========== Slide 2: Agenda ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide2, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Workshop Agenda", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))
    tf2 = add_textbox(slide2, Inches(0.8), Inches(1.8), Inches(11), Inches(5),
                      "9:00 AM - Introduction and Setup", font_size=18)
    add_bullet_text(tf2, "10:00 AM - Data Collection Strategies", level=0)
    add_bullet_text(tf2, "11:00 AM - Data Cleaning Techniques", level=0)
    add_bullet_text(tf2, "12:00 PM - Lunch Break", level=0)
    add_bullet_text(tf2, "1:00 PM - Exploratory Data Analysis", level=0)
    add_bullet_text(tf2, "2:30 PM - Visualization Best Practices", level=0)
    add_bullet_text(tf2, "3:30 PM - Hands-on Exercise", level=0)
    add_bullet_text(tf2, "4:30 PM - Wrap-Up and Q&A", level=0)

    # ========== Slide 3: Data Collection ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide3, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Data Collection Strategies", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))
    tf3 = add_textbox(slide3, Inches(0.8), Inches(1.8), Inches(5), Inches(5),
                      "Primary sources: surveys, interviews, sensor data", font_size=16)
    add_bullet_text(tf3, "Secondary sources: public datasets, APIs, web scraping", level=0, font_size=16)
    add_bullet_text(tf3, "Real-time streams: IoT devices, transaction logs", level=0, font_size=16)
    add_bullet_text(tf3, "Data quality checklist:", level=0, font_size=16)
    add_bullet_text(tf3, "Completeness - are all required fields present?", level=1, font_size=14)
    add_bullet_text(tf3, "Accuracy - does data reflect reality?", level=1, font_size=14)
    add_bullet_text(tf3, "Timeliness - is data current enough?", level=1, font_size=14)

    # ========== Slide 4: Data Cleaning ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide4, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Data Cleaning Techniques", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))
    tf4 = add_textbox(slide4, Inches(0.8), Inches(1.8), Inches(11), Inches(5),
                      "Handling Missing Values", font_size=20, bold=True)
    add_bullet_text(tf4, "Mean/median imputation for numerical data", level=0, font_size=16)
    add_bullet_text(tf4, "Mode imputation for categorical fields", level=0, font_size=16)
    add_bullet_text(tf4, "Forward-fill / backward-fill for time series", level=0, font_size=16)
    add_bullet_text(tf4, "", level=0, font_size=16)
    add_bullet_text(tf4, "Outlier Detection", level=0, font_size=20)
    add_bullet_text(tf4, "Z-score method (threshold > 3 standard deviations)", level=0, font_size=16)
    add_bullet_text(tf4, "IQR method (1.5 * interquartile range)", level=0, font_size=16)
    add_bullet_text(tf4, "Visual inspection with box plots", level=0, font_size=16)

    # ========== Slide 5: Exploratory Data Analysis ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide5, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Exploratory Data Analysis (EDA)", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))
    tf5 = add_textbox(slide5, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5),
                      "Descriptive statistics: mean, median, std, skewness", font_size=16)
    add_bullet_text(tf5, "Correlation analysis: Pearson, Spearman coefficients", level=0, font_size=16)
    add_bullet_text(tf5, "Distribution shape: histograms, density plots", level=0, font_size=16)
    add_bullet_text(tf5, "Group comparisons: t-tests, ANOVA", level=0, font_size=16)
    add_bullet_text(tf5, "Feature relationships: scatter matrix, heatmaps", level=0, font_size=16)

    # ========== Slide 6: Visualization Best Practices ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide6, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Visualization Best Practices", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))
    tf6 = add_textbox(slide6, Inches(0.8), Inches(1.8), Inches(11), Inches(5),
                      "Choose the right chart type for your data:", font_size=18, bold=True)
    add_bullet_text(tf6, "Bar charts: comparing categories", level=0, font_size=16)
    add_bullet_text(tf6, "Line charts: trends over time", level=0, font_size=16)
    add_bullet_text(tf6, "Scatter plots: relationship between two variables", level=0, font_size=16)
    add_bullet_text(tf6, "Heatmaps: multi-dimensional patterns", level=0, font_size=16)
    add_bullet_text(tf6, "Pie charts: use sparingly, only for proportions", level=0, font_size=16)
    add_bullet_text(tf6, "", level=0, font_size=16)
    add_bullet_text(tf6, "Design principles: clarity, minimal ink, appropriate color palettes", level=0, font_size=16)

    # ========== Slide 7: Background - Statistics Refresher ==========
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide7, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Background: Statistics Refresher", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))
    tf7 = add_textbox(slide7, Inches(0.8), Inches(1.8), Inches(11), Inches(5),
                      "Central Tendency Measures", font_size=20, bold=True)
    add_bullet_text(tf7, "Mean: arithmetic average, sensitive to outliers", level=0, font_size=16)
    add_bullet_text(tf7, "Median: middle value, robust to extreme values", level=0, font_size=16)
    add_bullet_text(tf7, "Mode: most frequent value, useful for categorical data", level=0, font_size=16)
    add_bullet_text(tf7, "", level=0, font_size=16)
    add_bullet_text(tf7, "Spread Measures", level=0, font_size=20)
    add_bullet_text(tf7, "Standard deviation: average distance from mean", level=0, font_size=16)
    add_bullet_text(tf7, "Variance: squared standard deviation", level=0, font_size=16)
    add_bullet_text(tf7, "Range: difference between max and min", level=0, font_size=16)
    add_bullet_text(tf7, "IQR: range of the middle 50% of data", level=0, font_size=16)

    # ========== Slide 8: Background - Python Libraries ==========
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide8, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Background: Essential Python Libraries", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))
    tf8 = add_textbox(slide8, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5),
                      "pandas - data manipulation and analysis", font_size=16)
    add_bullet_text(tf8, "DataFrame operations, groupby, merge, pivot", level=1, font_size=14)
    add_bullet_text(tf8, "numpy - numerical computing foundation", level=0, font_size=16)
    add_bullet_text(tf8, "Array operations, linear algebra, statistics", level=1, font_size=14)
    add_bullet_text(tf8, "matplotlib - core plotting library", level=0, font_size=16)
    add_bullet_text(tf8, "Highly customizable, publication-quality figures", level=1, font_size=14)
    add_bullet_text(tf8, "seaborn - statistical visualization", level=0, font_size=16)
    add_bullet_text(tf8, "Built on matplotlib, beautiful default themes", level=1, font_size=14)
    add_bullet_text(tf8, "scikit-learn - machine learning toolkit", level=0, font_size=16)
    add_bullet_text(tf8, "Classification, regression, clustering, preprocessing", level=1, font_size=14)

    # ========== Slide 9: Hands-on Exercise ==========
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide9, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Hands-on Exercise", font_size=32, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))
    tf9 = add_textbox(slide9, Inches(0.8), Inches(1.8), Inches(11), Inches(5),
                      "Exercise: Analyze the retail_sales_2024.csv dataset", font_size=18, bold=True)
    add_bullet_text(tf9, "Step 1: Load data and inspect first 10 rows", level=0, font_size=16)
    add_bullet_text(tf9, "Step 2: Calculate summary statistics per region", level=0, font_size=16)
    add_bullet_text(tf9, "Step 3: Identify and handle missing values", level=0, font_size=16)
    add_bullet_text(tf9, "Step 4: Create a bar chart of total sales by category", level=0, font_size=16)
    add_bullet_text(tf9, "Step 5: Build a correlation heatmap for numeric columns", level=0, font_size=16)
    add_bullet_text(tf9, "", level=0, font_size=16)
    add_bullet_text(tf9, "Time allotted: 60 minutes  |  Work in pairs  |  Ask questions anytime!", level=0, font_size=16)

    # ========== Slide 10: Conclusions & Next Steps ==========
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    fill10 = slide10.background.fill
    fill10.solid()
    fill10.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    add_textbox(slide10, Inches(1), Inches(0.8), Inches(11), Inches(1.2),
                "Conclusions & Next Steps", font_size=36, bold=True,
                alignment=PP_ALIGN.CENTER, color=RGBColor(0xFF, 0xFF, 0xFF))
    tf10 = add_textbox(slide10, Inches(1), Inches(2.5), Inches(11), Inches(4),
                       "Key Takeaways", font_size=22, bold=True,
                       color=RGBColor(0xFF, 0xFF, 0xFF))
    p = add_bullet_text(tf10, "Data quality is the foundation of any analysis", level=0, font_size=18)
    p.runs[0].font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    p = add_bullet_text(tf10, "EDA should precede any modeling effort", level=0, font_size=18)
    p.runs[0].font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    p = add_bullet_text(tf10, "Visualizations tell stories that numbers alone cannot", level=0, font_size=18)
    p.runs[0].font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    p = add_bullet_text(tf10, "", level=0, font_size=18)
    p = add_bullet_text(tf10, "Resources: workshop materials at analytics.techhub.io/2025", level=0, font_size=16)
    p.runs[0].font.color.rgb = RGBColor(0x99, 0xBB, 0xDD)
    p = add_bullet_text(tf10, "Contact: elena.vasquez@techhub.io", level=0, font_size=16)
    p.runs[0].font.color.rgb = RGBColor(0x99, 0xBB, 0xDD)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
