"""
Reward Script: Reorganize slides, apply Cube transition, add slide numbers
Task ID: impress_gf4_009
Domain: libreoffice_impress
Scoring:
  Component 1 (0.40): Slides 7-8 moved to positions 3-4 (background slides before main content)
  Component 2 (0.35): All 10 slides have Cube (prism) transition applied
  Component 3 (0.25): All slides have slide number placeholders
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_009'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have 10 slides
    num_slides = len(prs.slides)
    if num_slides != 10:
        print(f"PRECONDITION FAIL: Expected 10 slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide reordering — slides 7&8 moved to positions 3&4 (0.40 points)
    # After reordering, slide 3 should be "Background: Statistics Refresher"
    # and slide 4 should be "Background: Essential Python Libraries"
    try:
        def get_slide_title_text(slide):
            """Extract the first non-empty text from the slide for identification."""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            return t
            return ""

        slide3_title = get_slide_title_text(prs.slides[2])  # 0-indexed
        slide4_title = get_slide_title_text(prs.slides[3])

        slide3_ok = "statistics refresher" in slide3_title.lower()
        slide4_ok = "essential python" in slide4_title.lower() or "python libraries" in slide4_title.lower()

        if slide3_ok and slide4_ok:
            print(f"PASS: Component 1 — Slide 3 is '{slide3_title}', Slide 4 is '{slide4_title}' (0.40 pts)")
            total_score += 0.40
        elif slide3_ok or slide4_ok:
            print(f"PARTIAL: Component 1 — Slide 3 ok={slide3_ok} ('{slide3_title}'), Slide 4 ok={slide4_ok} ('{slide4_title}') (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 — Slide 3 is '{slide3_title}' (expected 'Background: Statistics Refresher'), Slide 4 is '{slide4_title}' (expected 'Background: Essential Python Libraries')")

        # Additionally verify that slide order integrity is maintained
        # Slide 5 should be Data Collection, Slide 10 should be Conclusions
        slide5_title = get_slide_title_text(prs.slides[4])
        slide10_title = get_slide_title_text(prs.slides[9])
        if "data collection" in slide5_title.lower():
            print(f"  INFO: Slide 5 correctly shows '{slide5_title}'")
        else:
            print(f"  INFO: Slide 5 shows '{slide5_title}' (expected Data Collection)")
        if "conclusions" in slide10_title.lower() or "next steps" in slide10_title.lower():
            print(f"  INFO: Slide 10 correctly shows '{slide10_title}'")
        else:
            print(f"  INFO: Slide 10 shows '{slide10_title}' (expected Conclusions)")

    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: All slides have Cube/prism transition (0.35 points)
    # LibreOffice Cube maps to OOXML p14:prism
    try:
        ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        ns_p14 = 'http://schemas.microsoft.com/office/powerpoint/2010/main'

        slides_with_transition = 0
        with zipfile.ZipFile(file_path, 'r') as zf:
            for i in range(1, num_slides + 1):
                try:
                    with zf.open(f'ppt/slides/slide{i}.xml') as f:
                        root = ET.parse(f).getroot()
                        # Look for transition element
                        tr = root.find(f'{{{ns_p}}}transition')
                        if tr is None:
                            # Also search deeper
                            for elem in root.iter():
                                tag = elem.tag
                                if 'transition' in tag.lower():
                                    tr = elem
                                    break

                        if tr is not None:
                            # Check for prism (Cube) child - could be in p14 namespace
                            cube_children = [c for c in tr
                                             if (c.tag.split('}')[1] if '}' in c.tag else c.tag).lower()
                                             in ('prism', 'cube')]
                            if len(cube_children) > 0:
                                slides_with_transition += 1
                            else:
                                child_tags = [c.tag.split('}')[1] if '}' in c.tag else c.tag for c in tr]
                                print(f"  INFO: Slide {i} has transition but not Cube/prism: {child_tags}")
                        else:
                            print(f"  INFO: Slide {i} has no transition")
                except KeyError:
                    print(f"  INFO: Slide {i} XML not found")

        if slides_with_transition == num_slides:
            print(f"PASS: Component 2 — All {num_slides} slides have Cube transition (0.35 pts)")
            total_score += 0.35
        elif slides_with_transition > 0:
            partial = 0.35 * (slides_with_transition / num_slides)
            print(f"PARTIAL: Component 2 — {slides_with_transition}/{num_slides} slides have Cube transition ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No slides have Cube transition")

    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide numbers present on all slides (0.25 points)
    # Check for slidenum field in XML (fld type="slidenum" or placeholder idx=12)
    try:
        slides_with_number = 0
        with zipfile.ZipFile(file_path, 'r') as zf:
            for i in range(1, num_slides + 1):
                try:
                    with zf.open(f'ppt/slides/slide{i}.xml') as f:
                        content = f.read().decode('utf-8')
                        # Check for slidenum field type or placeholder type sldNum
                        if 'slidenum' in content.lower() or 'sldNum' in content:
                            slides_with_number += 1
                        else:
                            print(f"  INFO: Slide {i} has no slide number reference")
                except KeyError:
                    print(f"  INFO: Slide {i} XML not found")

        if slides_with_number == num_slides:
            print(f"PASS: Component 3 — All {num_slides} slides have slide numbers (0.25 pts)")
            total_score += 0.25
        elif slides_with_number > 0:
            partial = 0.25 * (slides_with_number / num_slides)
            print(f"PARTIAL: Component 3 — {slides_with_number}/{num_slides} slides have slide numbers ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No slides have slide numbers")

    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
