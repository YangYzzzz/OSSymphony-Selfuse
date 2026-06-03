"""
Initial Setup: Create a presentation with 8 slides, slide 5 titled 'Gallery' with empty content area.
Task ID: impress_gf3_012
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_012'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_placeholder(slide, idx, text, font_size=Pt(18)):
    """Add text to a placeholder by index if it exists."""
    if idx in slide.placeholders:
        tf = slide.placeholders[idx].text_frame
        tf.paragraphs[0].text = text
        for run in tf.paragraphs[0].runs:
            run.font.size = font_size


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Annual Photo Collection 2025"
    slide1.placeholders[1].text = "Capturing Moments Across the Year"

    # --- Slide 2: Introduction ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Introduction"
    tf2 = slide2.placeholders[1].text_frame
    tf2.paragraphs[0].text = "This presentation showcases the best photographs collected throughout 2025."
    p2 = tf2.add_paragraph()
    p2.text = "From team outings to product launches, these images capture key moments in our company history."
    p3 = tf2.add_paragraph()
    p3.text = "Each section highlights a different theme or event category."

    # --- Slide 3: Timeline ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Timeline of Events"
    tf3 = slide3.placeholders[1].text_frame
    tf3.paragraphs[0].text = "Q1 (Jan-Mar): Winter Retreat & New Year Kickoff"
    for item in [
        "Q2 (Apr-Jun): Spring Product Launch & Team Building",
        "Q3 (Jul-Sep): Summer Conference & Client Visits",
        "Q4 (Oct-Dec): Holiday Celebration & Year-End Review",
    ]:
        p = tf3.add_paragraph()
        p.text = item

    # --- Slide 4: Team Photos ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Team Photos"
    tf4 = slide4.placeholders[1].text_frame
    tf4.paragraphs[0].text = "Engineering Department - 42 members across 6 teams"
    for item in [
        "Marketing & Design - Creative campaigns and brand management",
        "Operations & Logistics - Ensuring smooth day-to-day execution",
        "Executive Leadership - Strategic vision and company direction",
    ]:
        p = tf4.add_paragraph()
        p.text = item

    # --- Slide 5: Gallery (EMPTY - task target) ---
    # Use layout 5 (blank) and add only a title text box manually
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title as a text box at the top
    txBox = slide5.shapes.add_textbox(Cm(1), Cm(0.5), Cm(23), Cm(2))
    tf5 = txBox.text_frame
    tf5.paragraphs[0].text = "Gallery"
    tf5.paragraphs[0].alignment = PP_ALIGN.LEFT
    for run in tf5.paragraphs[0].runs:
        run.font.size = Pt(36)
        run.font.bold = True
    # NO other shapes - this is the empty content area for the task

    # --- Slide 6: Events ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Notable Events"
    tf6 = slide6.placeholders[1].text_frame
    tf6.paragraphs[0].text = "March 15 - Annual Tech Conference in San Francisco"
    for item in [
        "June 8 - Product v3.0 Launch Event at HQ",
        "September 22 - Client Appreciation Dinner, Tokyo Office",
        "December 5 - Year-End Gala & Awards Ceremony",
    ]:
        p = tf6.add_paragraph()
        p.text = item

    # --- Slide 7: Summary ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Summary & Next Steps"
    tf7 = slide7.placeholders[1].text_frame
    tf7.paragraphs[0].text = "Over 500 photographs were collected and curated this year."
    for item in [
        "Top 50 images selected for the annual calendar.",
        "Photo contest winners will be announced at the January meeting.",
        "Submissions for 2026 open on February 1st.",
    ]:
        p = tf7.add_paragraph()
        p.text = item

    # --- Slide 8: Thank You ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    txBox8 = slide8.shapes.add_textbox(Cm(5), Cm(7), Cm(15), Cm(4))
    tf8 = txBox8.text_frame
    tf8.paragraphs[0].text = "Thank You"
    tf8.paragraphs[0].alignment = PP_ALIGN.CENTER
    for run in tf8.paragraphs[0].runs:
        run.font.size = Pt(44)
        run.font.bold = True
    p8 = tf8.add_paragraph()
    p8.text = "Questions? Contact photos@company.com"
    p8.alignment = PP_ALIGN.CENTER

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
