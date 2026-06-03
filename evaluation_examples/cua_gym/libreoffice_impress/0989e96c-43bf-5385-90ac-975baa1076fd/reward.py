"""
Reward Script: 3x3 image grid of rectangles on slide 5
Task ID: impress_gf3_012
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30) - Exactly 9 rectangle shapes on slide 5
  Component 2 (0.25) - All rectangles are 6cm x 5cm
  Component 3 (0.25) - Grid positioning: 3 rows x 3 cols with 0.5cm gaps, starting at (1cm, 3cm)
  Component 4 (0.20) - Fill color #EEEEEE and visible border on all rectangles
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_012'

# 1 cm = 360000 EMU
CM = 360000
TOLERANCE = 0.05  # 5% relative tolerance for position/size checks


def approx_eq(val, expected, tol=TOLERANCE):
    """Check if val is approximately equal to expected (relative tolerance)."""
    if expected == 0:
        return abs(val) < CM * 0.1  # within 0.1 cm of zero
    return abs(val - expected) / abs(expected) <= tol


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # slide 5 (0-indexed)

    # Collect all AUTO_SHAPE rectangles on slide 5
    # (filter out placeholders and text boxes which are pre-existing)
    rectangles = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            rectangles.append(shape)

    print(f"Found {len(rectangles)} auto-shapes on slide 5")

    # Component 1: Exactly 9 rectangle shapes on slide 5 (0.30 points)
    try:
        if len(rectangles) == 9:
            print(f"PASS: Component 1 - Exactly 9 rectangle shapes found (0.30 pts)")
            total_score += 0.30
        elif len(rectangles) >= 7:
            partial = 0.15
            print(f"PARTIAL: Component 1 - Found {len(rectangles)} rectangles, expected 9 ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 - Found {len(rectangles)} rectangle shapes, expected 9")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    if len(rectangles) == 0:
        print("No rectangles found, cannot verify further components.")
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: All rectangles are 6cm wide x 5cm tall (0.25 points)
    try:
        expected_w = 6 * CM   # 2160000 EMU
        expected_h = 5 * CM   # 1800000 EMU
        correct_size_count = 0
        for rect in rectangles:
            if approx_eq(rect.width, expected_w) and approx_eq(rect.height, expected_h):
                correct_size_count += 1
            else:
                print(f"  Size mismatch: {rect.name} w={rect.width/CM:.2f}cm h={rect.height/CM:.2f}cm")

        if correct_size_count == len(rectangles) and len(rectangles) >= 9:
            print(f"PASS: Component 2 - All {correct_size_count} rectangles are 6cm x 5cm (0.25 pts)")
            total_score += 0.25
        elif correct_size_count >= 7:
            partial = 0.15
            print(f"PARTIAL: Component 2 - {correct_size_count}/{len(rectangles)} correct size ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 - Only {correct_size_count}/{len(rectangles)} rectangles have correct size")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Grid positioning - 3 rows x 3 cols, 0.5cm gaps, starting at (1cm, 3cm) (0.25 points)
    try:
        gap = 0.5 * CM   # 180000 EMU
        start_x = 1 * CM  # 360000 EMU
        start_y = 3 * CM  # 1080000 EMU
        rect_w = 6 * CM
        rect_h = 5 * CM

        # Expected positions for 3x3 grid
        expected_positions = []
        for row in range(3):
            for col in range(3):
                ex = start_x + col * (rect_w + gap)
                ey = start_y + row * (rect_h + gap)
                expected_positions.append((ex, ey))

        # Match each rectangle to an expected position
        matched = 0
        used_positions = set()
        for rect in rectangles:
            for idx, (ex, ey) in enumerate(expected_positions):
                if idx not in used_positions:
                    if approx_eq(rect.left, ex) and approx_eq(rect.top, ey):
                        matched += 1
                        used_positions.add(idx)
                        break
            else:
                print(f"  Position mismatch: {rect.name} at ({rect.left/CM:.2f}cm, {rect.top/CM:.2f}cm)")

        if matched >= 9:
            print(f"PASS: Component 3 - All 9 rectangles correctly positioned in 3x3 grid (0.25 pts)")
            total_score += 0.25
        elif matched >= 6:
            partial = 0.15
            print(f"PARTIAL: Component 3 - {matched}/9 rectangles correctly positioned ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 - Only {matched}/9 rectangles correctly positioned")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Fill color #EEEEEE and visible border on all rectangles (0.20 points)
    try:
        correct_fill_count = 0
        correct_border_count = 0
        for rect in rectangles:
            # Check fill
            try:
                if rect.fill.type is not None:
                    color_str = str(rect.fill.fore_color.rgb).upper()
                    if color_str == 'EEEEEE':
                        correct_fill_count += 1
                    else:
                        print(f"  Fill mismatch: {rect.name} has fill {color_str}")
                else:
                    print(f"  Fill mismatch: {rect.name} has no fill")
            except Exception:
                print(f"  Fill check error for {rect.name}")

            # Check border (line)
            try:
                line = rect.line
                if line.fill.type is not None or (line.width is not None and line.width > 0):
                    correct_border_count += 1
                else:
                    print(f"  Border mismatch: {rect.name} has no visible border")
            except Exception:
                print(f"  Border check error for {rect.name}")

        fill_pass = correct_fill_count >= 9
        border_pass = correct_border_count >= 9

        if fill_pass and border_pass:
            print(f"PASS: Component 4 - All rectangles have #EEEEEE fill and visible border (0.20 pts)")
            total_score += 0.20
        elif correct_fill_count >= 7 or correct_border_count >= 7:
            partial = 0.10
            print(f"PARTIAL: Component 4 - {correct_fill_count} correct fill, {correct_border_count} with border ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 - {correct_fill_count} correct fill, {correct_border_count} with border")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
