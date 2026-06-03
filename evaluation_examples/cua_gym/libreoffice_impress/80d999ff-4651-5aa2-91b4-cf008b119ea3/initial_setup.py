"""
Initial Setup: Blank presentation with default slide size for poster task
Task ID: impress_gf2_018
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Cm

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_018'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # Default slide size: 25.4 x 19.05 cm (standard 10x7.5 inches)
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(19.05)

    # Add a single blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only layout (blank-ish)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
