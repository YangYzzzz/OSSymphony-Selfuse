"""
Reward Script: University Research Poster - A0 Landscape with Column Layout
Task ID: impress_gf2_018
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Slide size is A0 landscape (118.9 x 84.1 cm)
  Component 2 (0.25): Title banner rectangle - full width, ~12cm tall, filled #003366
  Component 3 (0.20): Title text - 60pt bold white text inside the banner area
  Component 4 (0.30): Three column rectangles below banner with correct fills
"""

import os
from pptx import Presentation
from pptx.util import Emu, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_018'

# Tolerances
CM_TOLERANCE = 2.0  # cm tolerance for position/size checks
COLOR_TOLERANCE = 10  # RGB channel tolerance


def approx_cm(emu_val, expected_cm, tol=CM_TOLERANCE):
    """Check if EMU value is approximately equal to expected cm."""
    actual_cm = emu_val / 360000.0
    return abs(actual_cm - expected_cm) <= tol


def color_close(actual_rgb, expected_hex, tol=COLOR_TOLERANCE):
    """Check if an RGBColor is close to expected hex string."""
    if actual_rgb is None:
        return False
    actual_str = str(actual_rgb).upper()
    expected_hex = expected_hex.upper()
    # Parse both to ints
    try:
        ar = int(actual_str[0:2], 16)
        ag = int(actual_str[2:4], 16)
        ab = int(actual_str[4:6], 16)
        er = int(expected_hex[0:2], 16)
        eg = int(expected_hex[2:4], 16)
        eb = int(expected_hex[4:6], 16)
        return abs(ar - er) <= tol and abs(ag - eg) <= tol and abs(ab - eb) <= tol
    except (ValueError, IndexError):
        return False


def get_shape_fill_rgb(shape):
    """Get fill color of a shape, return hex string or None."""
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb).upper()
    except Exception:
        pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    if len(prs.slides) == 0:
        print("FAIL: No slides found")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]
    slide_width_cm = prs.slide_width / 360000.0
    slide_height_cm = prs.slide_height / 360000.0

    # Component 1: Slide size is A0 landscape (118.9 x 84.1 cm) — 0.25 points
    try:
        width_ok = abs(slide_width_cm - 118.9) <= 1.0
        height_ok = abs(slide_height_cm - 84.1) <= 1.0
        if width_ok and height_ok:
            print(f"PASS: Component 1 — Slide size is A0 landscape ({slide_width_cm:.1f} x {slide_height_cm:.1f} cm) (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — Expected ~118.9 x 84.1 cm, found {slide_width_cm:.1f} x {slide_height_cm:.1f} cm")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Title banner rectangle - full width, ~12cm tall, filled #003366 — 0.25 points
    try:
        banner_found = False
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or shape.shape_type == 1:
                fill_hex = get_shape_fill_rgb(shape)
                if fill_hex and color_close(RGBColor.from_string(fill_hex), "003366"):
                    # Check it's at top, full width, reasonable height
                    shape_top_cm = shape.top / 360000.0
                    shape_width_cm = shape.width / 360000.0
                    shape_height_cm = shape.height / 360000.0
                    top_ok = shape_top_cm <= 2.0  # near top
                    width_ok = shape_width_cm >= slide_width_cm * 0.9  # nearly full width
                    height_ok = 5.0 <= shape_height_cm <= 25.0  # reasonable banner height
                    if top_ok and width_ok and height_ok:
                        banner_found = True
                        print(f"PASS: Component 2 — Title banner found: {shape_width_cm:.1f}x{shape_height_cm:.1f}cm at top, fill={fill_hex} (0.25 pts)")
                        break
        if banner_found:
            total_score += 0.25
        else:
            print("FAIL: Component 2 — No full-width banner rectangle with #003366 fill found at top")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Title text - 60pt bold white text inside the banner area — 0.20 points
    try:
        title_text_found = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape_top_cm = shape.top / 360000.0
                # Title text should be in the banner area (top ~15cm)
                if shape_top_cm <= 15.0:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            text = (run.text or "").strip()
                            if len(text) < 3:
                                continue
                            # Check bold
                            is_bold = run.font.bold is True
                            # Check size ~60pt (762000 EMU = 60pt, allow some tolerance)
                            size_ok = False
                            if run.font.size is not None:
                                size_pt = run.font.size / 12700.0
                                size_ok = abs(size_pt - 60.0) <= 10.0
                            # Check white color
                            is_white = False
                            try:
                                if run.font.color.type is not None:
                                    rgb = str(run.font.color.rgb).upper()
                                    is_white = color_close(run.font.color.rgb, "FFFFFF")
                            except Exception:
                                pass

                            if is_bold and size_ok and is_white:
                                title_text_found = True
                                print(f"PASS: Component 3 — Title text found: '{text[:50]}...', bold={is_bold}, size={run.font.size/12700:.0f}pt, white (0.20 pts)")
                                break
                        if title_text_found:
                            break
            if title_text_found:
                break

        if title_text_found:
            total_score += 0.20
        else:
            print("FAIL: Component 3 — No 60pt bold white title text found in banner area")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Three column rectangles below banner with correct fills — 0.30 points
    # Left (#F0F4FF), Center (#FFFFFF), Right (#F0F4FF)
    try:
        # Find rectangles below the banner (top > 3cm) with solid fills
        column_rects = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or shape.shape_type == 1:
                shape_top_cm = shape.top / 360000.0
                shape_height_cm = shape.height / 360000.0
                fill_hex = get_shape_fill_rgb(shape)
                # Below the banner area and tall enough to be a column
                if shape_top_cm >= 3.0 and shape_height_cm >= 20.0 and fill_hex is not None:
                    column_rects.append({
                        'left_cm': shape.left / 360000.0,
                        'top_cm': shape_top_cm,
                        'width_cm': shape.width / 360000.0,
                        'height_cm': shape_height_cm,
                        'fill': fill_hex
                    })

        # Sort by left position
        column_rects.sort(key=lambda r: r['left_cm'])

        col_score = 0.0

        if len(column_rects) >= 3:
            # Check we have 3 columns
            col_score += 0.10
            print(f"  Found {len(column_rects)} column rectangles")

            # Check left column fill (#F0F4FF)
            if color_close(RGBColor.from_string(column_rects[0]['fill']), "F0F4FF"):
                col_score += 0.05
                print(f"  Left column fill OK: {column_rects[0]['fill']}")
            else:
                print(f"  Left column fill WRONG: expected ~F0F4FF, got {column_rects[0]['fill']}")

            # Check center column fill (#FFFFFF)
            if color_close(RGBColor.from_string(column_rects[1]['fill']), "FFFFFF"):
                col_score += 0.05
                print(f"  Center column fill OK: {column_rects[1]['fill']}")
            else:
                print(f"  Center column fill WRONG: expected ~FFFFFF, got {column_rects[1]['fill']}")

            # Check right column fill (#F0F4FF)
            if color_close(RGBColor.from_string(column_rects[2]['fill']), "F0F4FF"):
                col_score += 0.05
                print(f"  Right column fill OK: {column_rects[2]['fill']}")
            else:
                print(f"  Right column fill WRONG: expected ~F0F4FF, got {column_rects[2]['fill']}")

            # Check columns are roughly equal width (within 20% of each other)
            widths = [r['width_cm'] for r in column_rects[:3]]
            avg_width = sum(widths) / 3
            widths_similar = all(abs(w - avg_width) / avg_width <= 0.25 for w in widths)
            if widths_similar:
                col_score += 0.05
                print(f"  Column widths approximately equal: {[f'{w:.1f}' for w in widths]} cm")
            else:
                print(f"  Column widths not equal: {[f'{w:.1f}' for w in widths]} cm")

            if col_score > 0:
                print(f"PASS: Component 4 — Three-column layout verified ({col_score:.2f} pts)")
                total_score += col_score
        else:
            print(f"FAIL: Component 4 — Expected 3 column rectangles below banner, found {len(column_rects)}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
