"""
Reward Script: Product launch slide deck with custom color theme, transitions, notes, and exports
Task ID: impress_gf4_034
Domain: libreoffice_impress
Scoring:
  Component 1: Dissolve transitions on all 8 slides (0.25)
  Component 2: Speaker notes on all 8 slides, >= 2 sentences each (0.25)
  Component 3: Title text uses primary color #0EA5E9 (0.15)
  Component 4: Logo image present on slides (0.15)
  Component 5: Metadata title and author set correctly (0.10)
  Component 6: Desktop export files exist (PDF + PPTX) (0.10)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_034'

def verify_task(file_path):
    """Verify task completion with progressive scoring. Returns float 0.0-1.0."""
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    if num_slides != 8:
        print(f"PRECONDITION FAIL: Expected 8 slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Dissolve transitions on all 8 slides (0.25 points)
    try:
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        dissolve_count = 0
        with zipfile.ZipFile(file_path, 'r') as zf:
            for i in range(1, 9):
                with zf.open(f'ppt/slides/slide{i}.xml') as f:
                    root = ET.parse(f).getroot()
                    tr = root.find('.//p:transition', ns)
                    if tr is not None and tr.find('.//p:dissolve', ns) is not None:
                        dissolve_count += 1
        if dissolve_count == 8:
            print(f"PASS: Component 1 — All 8 slides have Dissolve transitions (0.25 pts)")
            total_score += 0.25
        elif dissolve_count > 0:
            partial = round(0.25 * dissolve_count / 8, 3)
            print(f"PARTIAL: Component 1 — {dissolve_count}/8 slides have Dissolve transitions ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — No slides have Dissolve transitions")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Speaker notes on all 8 slides, at least 2 sentences each (0.25 points)
    try:
        notes_pass_count = 0
        for i, slide in enumerate(prs.slides):
            try:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                # Count sentences: split by period, question mark, exclamation
                import re
                sentences = [s.strip() for s in re.split(r'[.!?]+', notes_text) if s.strip()]
                if len(sentences) >= 2:
                    notes_pass_count += 1
                else:
                    print(f"  Slide {i+1} notes too short: {len(sentences)} sentence(s)")
            except Exception:
                print(f"  Slide {i+1} has no notes")
        if notes_pass_count == 8:
            print(f"PASS: Component 2 — All 8 slides have speaker notes with >=2 sentences (0.25 pts)")
            total_score += 0.25
        elif notes_pass_count > 0:
            partial = round(0.25 * notes_pass_count / 8, 3)
            print(f"PARTIAL: Component 2 — {notes_pass_count}/8 slides have adequate notes ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No slides have adequate speaker notes")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Title text uses primary color #0EA5E9 (0.15 points)
    # In the golden file, the first text shape (title) on each slide has color 0EA5E9
    try:
        title_color_count = 0
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # First text shape with content is the title
                    txt = shape.text_frame.text.strip()
                    if txt:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                try:
                                    if run.font.color and run.font.color.type is not None:
                                        rgb_str = str(run.font.color.rgb).upper()
                                        if rgb_str == '0EA5E9':
                                            title_color_count += 1
                                except Exception:
                                    pass
                        break  # only check first text shape (title)
        if title_color_count >= 8:
            print(f"PASS: Component 3 — All 8 slide titles use primary color #0EA5E9 (0.15 pts)")
            total_score += 0.15
        elif title_color_count > 0:
            partial = round(0.15 * min(title_color_count, 8) / 8, 3)
            print(f"PARTIAL: Component 3 — {title_color_count}/8 titles use #0EA5E9 ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No slide titles use primary color #0EA5E9")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Logo image present on slides (0.15 points)
    # Golden file has an image on each slide (the logo)
    try:
        slides_with_image = 0
        for i, slide in enumerate(prs.slides):
            has_img = False
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    has_img = True
                    break
            if has_img:
                slides_with_image += 1
        if slides_with_image >= 8:
            print(f"PASS: Component 4 — All 8 slides have logo image (0.15 pts)")
            total_score += 0.15
        elif slides_with_image > 0:
            partial = round(0.15 * slides_with_image / 8, 3)
            print(f"PARTIAL: Component 4 — {slides_with_image}/8 slides have image ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — No slides have a logo image")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Metadata title='Product X Launch' and creator='Marketing Team' (0.10 points)
    try:
        metadata_score = 0.0
        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('docProps/core.xml') as f:
                core_xml = f.read().decode()
                root = ET.fromstring(core_xml)
                # Namespaces
                dc_ns = 'http://purl.org/dc/elements/1.1/'
                title_el = root.find(f'{{{dc_ns}}}title')
                creator_el = root.find(f'{{{dc_ns}}}creator')

                title_val = title_el.text.strip() if title_el is not None and title_el.text else ''
                creator_val = creator_el.text.strip() if creator_el is not None and creator_el.text else ''

                if title_val == 'Product X Launch':
                    metadata_score += 0.05
                    print(f"  Metadata title: '{title_val}' — correct")
                else:
                    print(f"  Metadata title: '{title_val}' — expected 'Product X Launch'")

                if creator_val == 'Marketing Team':
                    metadata_score += 0.05
                    print(f"  Metadata creator: '{creator_val}' — correct")
                else:
                    print(f"  Metadata creator: '{creator_val}' — expected 'Marketing Team'")

        if metadata_score > 0:
            print(f"PASS: Component 5 — Metadata checks ({metadata_score} pts)")
            total_score += metadata_score
        else:
            print(f"FAIL: Component 5 — No metadata matches")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Desktop export files (Product_X_Launch.pdf and Product_X_Launch.pptx) (0.10 points)
    try:
        export_score = 0.0
        pdf_path = '/home/user/Desktop/Product_X_Launch.pdf'
        pptx_path = '/home/user/Desktop/Product_X_Launch.pptx'

        if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 1000:
            export_score += 0.05
            print(f"  PDF export found: {os.path.getsize(pdf_path)} bytes")
        else:
            print(f"  PDF export missing or too small")

        if os.path.exists(pptx_path) and os.path.getsize(pptx_path) > 1000:
            export_score += 0.05
            print(f"  PPTX export found: {os.path.getsize(pptx_path)} bytes")
        else:
            print(f"  PPTX export missing or too small")

        if export_score > 0:
            print(f"PASS: Component 6 — Desktop exports ({export_score} pts)")
            total_score += export_score
        else:
            print(f"FAIL: Component 6 — No valid export files on Desktop")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Main entry point
persist_app_state()
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
