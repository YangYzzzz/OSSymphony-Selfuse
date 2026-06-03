"""
Initial Setup: Product launch slide deck - blank 8-slide presentation
Task ID: impress_gf4_034
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_034'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
LOGO_PATH = f'{WORKDIR}/Desktop/logo.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_logo():
    """Create a simple company logo PNG on the Desktop."""
    img = Image.new('RGBA', (200, 80), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    # Draw a simple logo shape - blue circle with text
    draw.ellipse([10, 10, 70, 70], fill=(14, 165, 233, 255))
    draw.text((80, 25), "ACME", fill=(30, 30, 30, 255))
    os.makedirs(os.path.dirname(LOGO_PATH), exist_ok=True)
    img.save(LOGO_PATH)
    print(f'Logo created: {LOGO_PATH}')


def create_initial():
    prs = Presentation()

    slide_titles = [
        'Introduction',
        'The Problem',
        'Our Solution',
        'Key Features',
        'Technical Specs',
        'Pricing',
        'Timeline',
        'Call to Action',
    ]

    for i, title in enumerate(slide_titles):
        if i == 0:
            # Title slide layout
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = "Product X Launch Event"
        else:
            # Title + Content layout
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = f"Content for {title}"

    # No transitions, no speaker notes, no custom colors, no logo on master,
    # no footer, no metadata changes - this is the blank initial state.

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


# Create logo file first, then the presentation
create_logo()
create_initial()
