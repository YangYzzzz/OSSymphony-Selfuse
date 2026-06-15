"""
Reward Script: Interactive presentation setup verification
Task ID: impress_tm_095
Domain: libreoffice_impress
Scoring:
  Component 1: Video embedded on slide 3         (0.20)
  Component 2: Video auto-play on slide 3        (0.10)
  Component 3: Audio embedded on slide 5         (0.20)
  Component 4: Audio on-click trigger on slide 5 (0.10)
  Component 5: Notes on all 7 slides             (0.25)
  Component 6: Custom show 'Short Version' 1,3,5,7 (0.15)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_095'

NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_RELS = 'http://schemas.openxmlformats.org/package/2006/relationships'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: file must be a valid pptx (zip)
    try:
        zf = zipfile.ZipFile(file_path, 'r')
    except Exception as e:
        print(f"CRITICAL: Cannot open pptx {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        # Precondition: must have 7 slides
        from pptx import Presentation
        prs = Presentation(file_path)
        if len(prs.slides) != 7:
            print(f"PRECONDITION FAIL: Expected 7 slides, found {len(prs.slides)}")
            print("REWARD: 0.0")
            return 0.0
    except Exception as e:
        print(f"CRITICAL: Cannot load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # -------------------------------------------------------
    # Component 1: Video embedded on slide 3 (0.20 points)
    # Check that slide 3 has a video relationship in its rels
    # -------------------------------------------------------
    try:
        has_video = False
        try:
            with zf.open('ppt/slides/_rels/slide3.xml.rels') as f:
                rels_content = f.read().decode()
            rels_root = ET.fromstring(rels_content)
            for rel in rels_root.findall(f'{{{NS_RELS}}}Relationship'):
                rel_type = rel.get('Type', '')
                if 'video' in rel_type.lower():
                    has_video = True
                    break
        except KeyError:
            pass

        if has_video:
            print(f"PASS: Component 1 - Video relationship found on slide 3 (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 - No video relationship found on slide 3")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # -------------------------------------------------------
    # Component 2: Video auto-play on slide 3 (0.10 points)
    # Auto-play means timing has delay=0 condition (NOT onClick)
    # The golden has: cond delay="0" at the sequence start level,
    # with no onClick gating the media playFrom command
    # -------------------------------------------------------
    try:
        auto_play = False
        with zf.open('ppt/slides/slide3.xml') as f:
            slide3_content = f.read().decode()
        slide3_root = ET.fromstring(slide3_content)

        # Look for timing element containing cmd playFrom
        has_timing = False
        has_onclick_gate = False
        has_play_cmd = False

        for el in slide3_root.iter():
            tag = el.tag.split('}')[-1] if '}' in el.tag else el.tag
            if tag == 'timing':
                has_timing = True
            if tag == 'cmd' and 'playFrom' in el.get('cmd', ''):
                has_play_cmd = True
            # Check if there's an onClick condition gating the play
            if tag == 'cond' and el.get('evt') == 'onClick':
                has_onclick_gate = True

        # Auto-play: has timing + play command + no onClick gate
        if has_timing and has_play_cmd and not has_onclick_gate:
            auto_play = True

        if auto_play:
            print(f"PASS: Component 2 - Video on slide 3 is auto-play (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 - Video on slide 3 is not auto-play "
                  f"(timing={has_timing}, playCmd={has_play_cmd}, onClick={has_onclick_gate})")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # -------------------------------------------------------
    # Component 3: Audio embedded on slide 5 (0.20 points)
    # Check that slide 5 has an audio relationship in its rels
    # -------------------------------------------------------
    try:
        has_audio = False
        try:
            with zf.open('ppt/slides/_rels/slide5.xml.rels') as f:
                rels_content = f.read().decode()
            rels_root = ET.fromstring(rels_content)
            for rel in rels_root.findall(f'{{{NS_RELS}}}Relationship'):
                rel_type = rel.get('Type', '')
                if 'audio' in rel_type.lower():
                    has_audio = True
                    break
        except KeyError:
            pass

        if has_audio:
            print(f"PASS: Component 3 - Audio relationship found on slide 5 (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 3 - No audio relationship found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # -------------------------------------------------------
    # Component 4: Audio on-click trigger on slide 5 (0.10 points)
    # On-click means the timing has an onClick condition
    # -------------------------------------------------------
    try:
        on_click = False
        with zf.open('ppt/slides/slide5.xml') as f:
            slide5_content = f.read().decode()
        slide5_root = ET.fromstring(slide5_content)

        has_timing = False
        has_play_cmd = False
        has_onclick = False

        for el in slide5_root.iter():
            tag = el.tag.split('}')[-1] if '}' in el.tag else el.tag
            if tag == 'timing':
                has_timing = True
            if tag == 'cmd' and 'playFrom' in el.get('cmd', ''):
                has_play_cmd = True
            if tag == 'cond' and el.get('evt') == 'onClick':
                has_onclick = True

        if has_timing and has_play_cmd and has_onclick:
            on_click = True

        if on_click:
            print(f"PASS: Component 4 - Audio on slide 5 triggers on click (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 4 - Audio on slide 5 not on-click "
                  f"(timing={has_timing}, playCmd={has_play_cmd}, onClick={has_onclick})")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # -------------------------------------------------------
    # Component 5: Notes on all 7 slides (0.25 points)
    # Each slide must have a notesSlide with non-empty text
    # Score progressively: 0.25 * (slides_with_notes / 7)
    # -------------------------------------------------------
    try:
        notes_count = 0
        for si in range(1, 8):
            try:
                note_path = f'ppt/notesSlides/notesSlide{si}.xml'
                with zf.open(note_path) as f:
                    note_content = f.read().decode()
                note_root = ET.fromstring(note_content)
                # Extract all text from the notes
                texts = []
                for t_el in note_root.iter(f'{{{NS_A}}}t'):
                    if t_el.text and t_el.text.strip():
                        texts.append(t_el.text.strip())
                # Filter out just slide number placeholders
                meaningful = [t for t in texts if len(t) > 5]
                if meaningful:
                    notes_count += 1
                    print(f"  Slide {si} notes: present ({len(meaningful)} text segments)")
                else:
                    print(f"  Slide {si} notes: empty or trivial")
            except KeyError:
                print(f"  Slide {si} notes: no notesSlide file")

        if notes_count == 7:
            print(f"PASS: Component 5 - All 7 slides have notes (0.25 pts)")
            total_score += 0.25
        elif notes_count > 0:
            partial = round(0.25 * (notes_count / 7), 4)
            print(f"PARTIAL: Component 5 - {notes_count}/7 slides have notes ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 - No slides have notes")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    # -------------------------------------------------------
    # Component 6: Custom slideshow 'Short Version' with slides 1,3,5,7 (0.15 points)
    # Check presentation.xml for custShowLst with correct name and slide refs
    # -------------------------------------------------------
    try:
        with zf.open('ppt/presentation.xml') as f:
            pres_content = f.read().decode()
        pres_root = ET.fromstring(pres_content)

        # Build rId-to-slide-index mapping from sldIdLst
        slide_rids = []
        for sld_id in pres_root.findall(f'.//{{{NS_P}}}sldId'):
            rid = sld_id.get(f'{{{NS_R}}}id')
            if rid:
                slide_rids.append(rid)
        # slide_rids[0] = slide 1, slide_rids[1] = slide 2, etc.

        # Find custom show named 'Short Version'
        found_show = False
        correct_slides = False

        for cust_show in pres_root.findall(f'.//{{{NS_P}}}custShow'):
            name = cust_show.get('name', '')
            if name == 'Short Version':
                found_show = True
                # Get the slide rIds in this custom show
                show_rids = []
                for sld in cust_show.findall(f'.//{{{NS_P}}}sld'):
                    rid = sld.get(f'{{{NS_R}}}id')
                    if rid:
                        show_rids.append(rid)

                # Map rIds to 1-based slide indices
                show_indices = []
                for rid in show_rids:
                    if rid in slide_rids:
                        show_indices.append(slide_rids.index(rid) + 1)

                print(f"  Custom show 'Short Version' slides: {show_indices}")
                if show_indices == [1, 3, 5, 7]:
                    correct_slides = True

        if found_show and correct_slides:
            print(f"PASS: Component 6 - Custom show 'Short Version' with slides 1,3,5,7 (0.15 pts)")
            total_score += 0.15
        elif found_show:
            print(f"FAIL: Component 6 - Custom show found but wrong slides")
        else:
            print(f"FAIL: Component 6 - No custom show named 'Short Version'")
    except Exception as e:
        print(f"ERROR: Component 6 - {e}")

    zf.close()

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
