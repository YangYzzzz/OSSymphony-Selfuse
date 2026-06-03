"""
Initial Setup: Build a 7-slide presentation with no media, no notes, no custom shows.
Task ID: impress_tm_095
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_095'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=None, alignment=None):
    """Helper to add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16):
    """Helper to add a bulleted text list."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(font_size)
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ─── Slide 1: Title Slide ───
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x2A, 0x4A)
    add_text_box(slide1, 1.5, 1.0, 10, 1.5,
                 "Horizon Technologies", font_size=44, bold=True,
                 color=(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, 1.5, 2.8, 10, 1.0,
                 "Annual Product Strategy Review 2025", font_size=24,
                 color=(0xA0, 0xC4, 0xFF), alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, 1.5, 4.2, 10, 0.8,
                 "Presented by Elena Rodriguez, VP of Product", font_size=18,
                 color=(0xCC, 0xCC, 0xCC), alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, 1.5, 5.2, 10, 0.5,
                 "March 15, 2025  |  San Francisco HQ", font_size=14,
                 color=(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)

    # ─── Slide 2: Agenda ───
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide2, 0.8, 0.5, 11, 1.0,
                 "Agenda", font_size=36, bold=True,
                 color=(0x1B, 0x2A, 0x4A))
    agenda_items = [
        "1. Company Performance Overview & Key Metrics",
        "2. Product Demo: Next-Gen Platform (Video)",
        "3. Market Analysis & Competitive Landscape",
        "4. Customer Success Stories & Case Studies",
        "5. Sound Design Preview for Brand Refresh",
        "6. Q4 Roadmap & Strategic Priorities",
        "7. Q&A and Closing Remarks"
    ]
    add_bullet_list(slide2, 0.8, 1.8, 11, 5.0, agenda_items, font_size=20)

    # ─── Slide 3: Product Demo (where video will go) ───
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide3, 0.8, 0.5, 11, 1.0,
                 "Product Demo: Next-Gen Platform", font_size=36, bold=True,
                 color=(0x1B, 0x2A, 0x4A))
    add_text_box(slide3, 1.0, 2.0, 10, 1.0,
                 "Our latest platform showcases AI-driven analytics, "
                 "real-time collaboration, and enterprise-grade security.",
                 font_size=18, color=(0x33, 0x33, 0x33))
    # Placeholder box indicating where the demo video should be inserted
    placeholder = slide3.shapes.add_shape(
        1, Inches(2.5), Inches(3.2), Inches(8), Inches(3.5))  # Rectangle
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    tf = placeholder.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[Insert Product Demo Video Here]"
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    # ─── Slide 4: Market Analysis ───
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide4, 0.8, 0.5, 11, 1.0,
                 "Market Analysis & Competitive Landscape", font_size=36,
                 bold=True, color=(0x1B, 0x2A, 0x4A))
    market_items = [
        "Total Addressable Market: $4.2B (growing 18% YoY)",
        "Current market share: 12.3% (up from 8.7% in 2024)",
        "Key competitors: TechStar Corp, DataBridge Inc, CloudNova",
        "Our differentiator: AI-first architecture with 99.99% uptime",
        "Customer retention rate: 94.6% (industry avg: 82%)",
        "Net Promoter Score: 72 (industry avg: 45)"
    ]
    add_bullet_list(slide4, 0.8, 1.8, 11, 5.0, market_items, font_size=18)

    # ─── Slide 5: Sound Design (where audio will go) ───
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide5, 0.8, 0.5, 11, 1.0,
                 "Brand Refresh: Sound Design Preview", font_size=36,
                 bold=True, color=(0x1B, 0x2A, 0x4A))
    add_text_box(slide5, 0.8, 1.8, 11, 1.5,
                 "Our new brand sound identity has been crafted by award-winning "
                 "composer Yuki Tanaka. The signature audio cue will be used "
                 "across all customer touchpoints — from app notifications to "
                 "conference presentations.",
                 font_size=18, color=(0x33, 0x33, 0x33))
    sound_items = [
        "Signature brand chime — 3-second audio motif",
        "Notification sounds — subtle, professional tones",
        "Presentation intro/outro — cinematic quality",
        "Hold music — calming, on-brand ambient tracks"
    ]
    add_bullet_list(slide5, 0.8, 3.5, 11, 3.0, sound_items, font_size=18)

    # ─── Slide 6: Q4 Roadmap ───
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide6, 0.8, 0.5, 11, 1.0,
                 "Q4 2025 Roadmap & Strategic Priorities", font_size=36,
                 bold=True, color=(0x1B, 0x2A, 0x4A))
    roadmap_items = [
        "October: Launch AI Assistant v2.0 with multi-language support",
        "October: Complete SOC 2 Type II certification",
        "November: Release enterprise API gateway",
        "November: Expand APAC data center presence (Tokyo, Singapore)",
        "December: Beta launch of real-time collaboration features",
        "December: Annual customer summit in New York City"
    ]
    add_bullet_list(slide6, 0.8, 1.8, 11, 5.0, roadmap_items, font_size=18)

    # ─── Slide 7: Closing / Q&A ───
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide7.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x2A, 0x4A)
    add_text_box(slide7, 1.5, 2.0, 10, 1.5,
                 "Thank You", font_size=48, bold=True,
                 color=(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_text_box(slide7, 1.5, 3.8, 10, 1.0,
                 "Questions & Discussion", font_size=28,
                 color=(0xA0, 0xC4, 0xFF), alignment=PP_ALIGN.CENTER)
    add_text_box(slide7, 1.5, 5.2, 10, 0.8,
                 "elena.rodriguez@horizontech.com  |  @horizonTech",
                 font_size=16, color=(0xCC, 0xCC, 0xCC),
                 alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # ─── Create media files on VM ───
    os.makedirs(f'{WORKDIR}/videos', exist_ok=True)
    os.makedirs(f'{WORKDIR}/sounds', exist_ok=True)

    # Create a minimal valid MP4 file (ftyp box only is enough as a placeholder)
    _create_minimal_mp4(f'{WORKDIR}/videos/intro.mp4')
    _create_minimal_wav(f'{WORKDIR}/sounds/effect.wav')

    # ─── GUI-ready startup ───
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


def _create_minimal_mp4(path):
    """Create a minimal valid MP4 file using ffmpeg if available, else write a tiny stub."""
    import shutil
    if shutil.which('ffmpeg'):
        subprocess.run([
            'ffmpeg', '-y', '-f', 'lavfi', '-i',
            'color=c=black:s=320x240:d=3',
            '-f', 'lavfi', '-i', 'anullsrc=r=44100:cl=mono',
            '-t', '3', '-c:v', 'libx264', '-c:a', 'aac',
            '-shortest', path
        ], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=30)
        print(f'Created video: {path}')
    else:
        # Write minimal MP4 ftyp header
        import struct
        with open(path, 'wb') as f:
            ftyp = b'ftypisom' + b'\x00' * 4
            f.write(struct.pack('>I', len(ftyp) + 8) + b'ftyp' + ftyp[:4])
        print(f'Created stub video: {path}')


def _create_minimal_wav(path):
    """Create a minimal valid WAV file (1 second of silence)."""
    import struct
    sample_rate = 44100
    num_samples = sample_rate  # 1 second
    data_size = num_samples * 2  # 16-bit mono
    with open(path, 'wb') as f:
        # RIFF header
        f.write(b'RIFF')
        f.write(struct.pack('<I', 36 + data_size))
        f.write(b'WAVE')
        # fmt chunk
        f.write(b'fmt ')
        f.write(struct.pack('<I', 16))        # chunk size
        f.write(struct.pack('<H', 1))         # PCM
        f.write(struct.pack('<H', 1))         # mono
        f.write(struct.pack('<I', sample_rate))
        f.write(struct.pack('<I', sample_rate * 2))  # byte rate
        f.write(struct.pack('<H', 2))         # block align
        f.write(struct.pack('<H', 16))        # bits per sample
        # data chunk
        f.write(b'data')
        f.write(struct.pack('<I', data_size))
        f.write(b'\x00' * data_size)
    print(f'Created audio: {path}')


create_initial()
