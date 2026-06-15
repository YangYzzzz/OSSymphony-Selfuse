"""
Initial Setup: Physics Demo presentation with 9 slides.
Task ID: impress_stu_027
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_027'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Helper to add a title+content slide
    def add_titled_slide(title_text, body_text, layout_idx=1):
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        slide.shapes.title.text = title_text
        if layout_idx == 1 and len(slide.placeholders) > 1:
            slide.placeholders[1].text = body_text
        else:
            # Add a text box for the body
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = body_text
            run = p.runs[0]
            run.font.size = Pt(18)
        return slide

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Physics Demonstrations"
    slide1.placeholders[1].text = "Exploring Fundamental Concepts Through Experiments\nDr. Elena Vasquez | Department of Physics"

    # Slide 2: Outline
    add_titled_slide(
        "Course Overview",
        "This presentation covers key physics demonstrations including:\n"
        "- Newton's Laws of Motion\n"
        "- Conservation of Energy\n"
        "- Wave Phenomena\n"
        "- Electromagnetic Induction\n"
        "- Fluid Dynamics"
    )

    # Slide 3: Newton's Laws
    add_titled_slide(
        "Newton's Laws of Motion",
        "First Law: An object at rest stays at rest unless acted upon by an external force.\n"
        "Second Law: F = ma — Force equals mass times acceleration.\n"
        "Third Law: For every action, there is an equal and opposite reaction.\n\n"
        "Demonstration: Air track with frictionless gliders showing conservation of momentum."
    )

    # Slide 4: Conservation of Energy
    add_titled_slide(
        "Conservation of Energy",
        "Energy cannot be created or destroyed, only transformed.\n\n"
        "Kinetic Energy: KE = ½mv²\n"
        "Potential Energy: PE = mgh\n\n"
        "Demonstration: Pendulum swing showing KE ↔ PE conversion.\n"
        "Measured swing heights confirm energy conservation within 2% error margin."
    )

    # Slide 5: Wave Phenomena
    add_titled_slide(
        "Wave Phenomena",
        "Waves transfer energy without transferring matter.\n\n"
        "Key Properties:\n"
        "- Wavelength (λ): Distance between successive crests\n"
        "- Frequency (f): Number of oscillations per second\n"
        "- Velocity: v = λf\n\n"
        "Demonstration: Standing waves on a vibrating string at resonant frequencies."
    )

    # Slide 6: Electromagnetic Induction
    add_titled_slide(
        "Electromagnetic Induction",
        "Faraday's Law: A changing magnetic flux induces an electromotive force (EMF).\n\n"
        "ε = -dΦ/dt\n\n"
        "Applications: Electric generators, transformers, wireless charging.\n"
        "Demonstration: Moving a magnet through a coil to light an LED."
    )

    # Slide 7: Video Demonstration — NO hyperlink in initial state
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Video Demonstration"
    if len(slide7.placeholders) > 1:
        slide7.placeholders[1].text = (
            "In this section, we will review a recorded experiment that illustrates "
            "the principles discussed in the previous slides.\n\n"
            "The experiment demonstrates the relationship between applied force "
            "and acceleration using a low-friction cart on a calibrated track.\n\n"
            "Please pay attention to the measurement techniques and error analysis."
        )

    # Slide 8: Data Analysis
    add_titled_slide(
        "Experimental Data Analysis",
        "Trial 1: Force = 2.0 N, Acceleration = 1.98 m/s²\n"
        "Trial 2: Force = 4.0 N, Acceleration = 3.95 m/s²\n"
        "Trial 3: Force = 6.0 N, Acceleration = 5.91 m/s²\n"
        "Trial 4: Force = 8.0 N, Acceleration = 7.88 m/s²\n\n"
        "Linear regression: a = 0.986F + 0.012\n"
        "R² = 0.9997 — excellent agreement with Newton's Second Law."
    )

    # Slide 9: Conclusion
    add_titled_slide(
        "Conclusions & Next Steps",
        "Key Takeaways:\n"
        "- Classical mechanics principles verified experimentally\n"
        "- Measurement precision within 2% of theoretical values\n"
        "- Wave phenomena successfully demonstrated at multiple frequencies\n\n"
        "Next Session: Thermodynamics — Heat transfer and entropy.\n"
        "Office Hours: Tuesdays 2-4 PM, Room 312 Kline Physics Building."
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
