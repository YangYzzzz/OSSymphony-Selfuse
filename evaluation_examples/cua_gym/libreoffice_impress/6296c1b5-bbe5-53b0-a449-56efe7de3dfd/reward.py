"""
Reward Script: Insert hyperlink on slide 7 with styled text
Task ID: impress_stu_027
Domain: libreoffice_impress
Scoring:
  Component 1: Hyperlink text 'Watch the experiment video' exists on slide 7 (0.25)
  Component 2: Hyperlink URL is https://www.youtube.com/watch?v=example123 (0.30)
  Component 3: Font size is 16pt (0.15)
  Component 4: Font color is #0066CC (0.15)
  Component 5: Text is underlined (0.15)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_027'


def persist_app_state(domain: str):
    """Save any unsaved GUI edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.oxml.ns import qn
    except ImportError as e:
        print(f"CRITICAL: Missing library: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 7 slides
    if len(prs.slides) < 7:
        print(f"PRECONDITION FAIL: Need at least 7 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[6]  # 0-indexed, slide 7

    # Search all shapes on slide 7 for the hyperlink text and properties
    target_text = "Watch the experiment video"
    found_run = None
    found_hlink_url = None

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text and target_text.lower() in run.text.strip().lower():
                    found_run = run
                    # Check for hyperlink via XML (hlinkClick is inside a:rPr)
                    rPr = run._r.find(qn('a:rPr'))
                    hlinkClick = None
                    if rPr is not None:
                        hlinkClick = rPr.find(qn('a:hlinkClick'))
                    if hlinkClick is not None:
                        rId = hlinkClick.get(qn('r:id'))
                        if rId:
                            try:
                                rel = slide.part.rels[rId]
                                found_hlink_url = rel.target_ref
                            except Exception:
                                found_hlink_url = None
                    break
            if found_run is not None:
                break
        if found_run is not None:
            break

    # Component 1: Hyperlink text exists on slide 7 (0.25 points)
    try:
        if found_run is not None:
            print(f"PASS: Component 1 - Text '{found_run.text}' found on slide 7 (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 - Text '{target_text}' not found on slide 7")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Hyperlink URL is correct (0.30 points)
    expected_url = "https://www.youtube.com/watch?v=example123"
    try:
        if found_hlink_url is not None and found_hlink_url == expected_url:
            print(f"PASS: Component 2 - Hyperlink URL matches '{expected_url}' (0.30 pts)")
            total_score += 0.30
        elif found_hlink_url is not None:
            print(f"FAIL: Component 2 - Hyperlink URL is '{found_hlink_url}', expected '{expected_url}'")
        else:
            print(f"FAIL: Component 2 - No hyperlink found on the target text")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Font size is 16pt (0.15 points)
    try:
        if found_run is not None and found_run.font.size is not None:
            actual_size_pt = found_run.font.size.pt
            if abs(actual_size_pt - 16.0) < 0.5:
                print(f"PASS: Component 3 - Font size is {actual_size_pt}pt (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 3 - Font size is {actual_size_pt}pt, expected 16pt")
        else:
            print(f"FAIL: Component 3 - Font size not set or text not found")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Font color is #0066CC (0.15 points)
    try:
        if found_run is not None:
            color_rgb = None
            try:
                if found_run.font.color.type is not None:
                    color_rgb = str(found_run.font.color.rgb)
            except Exception:
                pass
            # Also check via XML as fallback
            if color_rgb is None:
                solidFill = found_run._r.find('.//' + qn('a:srgbClr'))
                if solidFill is not None:
                    color_rgb = solidFill.get('val')
            if color_rgb is not None and color_rgb.upper() == "0066CC":
                print(f"PASS: Component 4 - Font color is #{color_rgb} (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 4 - Font color is '{color_rgb}', expected '0066CC'")
        else:
            print(f"FAIL: Component 4 - Text not found, cannot check color")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Text is underlined (0.15 points)
    try:
        if found_run is not None:
            is_underlined = found_run.font.underline
            # Also check XML attribute 'u' for underline as fallback
            if is_underlined is None:
                rPr_elem = found_run._r.find(qn('a:rPr'))
                if rPr_elem is not None:
                    u_val = rPr_elem.get('u')
                    is_underlined = (u_val is not None and u_val != 'none')
            if is_underlined:
                print(f"PASS: Component 5 - Text is underlined (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 5 - Text is not underlined (underline={is_underlined})")
        else:
            print(f"FAIL: Component 5 - Text not found, cannot check underline")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
