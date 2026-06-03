"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 7, there's a chart and I'd like to make sure whoever reviews it checks the numbers carefully. How can I add a comment right next to the chart saying 'Check data accuracy' so it's harder to miss?
Generated: 2025-08-07 11:06:11
Status: success
Model: o4-mini
Total Steps: 1
"""

import os
from pptx import Presentation

def verify_task(file_path):
    print("Starting verification...")
    score = 0.0
    max_score = 1.0

    # Requirement 1: File existence (0.2)
    if os.path.exists(file_path):
        print(f"✓ File exists: {file_path} (0.2)")
        score += 0.2
    else:
        print(f"✗ File not found: {file_path}")
        final_score = min(score, max_score)
        print(f"REWARD: {final_score}")
        return final_score

    # Requirement 2: Load presentation
    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded, slides: {len(prs.slides)}")
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        final_score = min(score, max_score)
        print(f"REWARD: {final_score}")
        return final_score

    # Verify at least 7 slides exist
    if len(prs.slides) < 7:
        print(f"✗ Only {len(prs.slides)} slides, expected at least 7")
        final_score = min(score, max_score)
        print(f"REWARD: {final_score}")
        return final_score
    else:
        print("✓ Contains at least 7 slides")

    # Focus on slide 7 (index 6)
    slide = prs.slides[6]

    # Requirement 3: Chart presence on slide 7 (0.3)
    chart_shapes = [sh for sh in slide.shapes if hasattr(sh, "chart")]
    if chart_shapes:
        chart = chart_shapes[0]
        print("✓ Chart found on slide 7 (0.3)")
        score += 0.3
    else:
        print("✗ No chart found on slide 7")
        final_score = min(score, max_score)
        print(f"REWARD: {final_score}")
        return final_score

    # Requirement 4: Presence of comment text 'Check data accuracy' (0.3)
    comment_shapes = [sh for sh in slide.shapes if hasattr(sh, "text_frame") and "check data accuracy" in sh.text.strip().lower()]
    if comment_shapes:
        comment = comment_shapes[0]
        print("✓ Found comment text 'Check data accuracy' (0.3)")
        score += 0.3
    else:
        print("✗ No comment text 'Check data accuracy' found on slide 7")
        final_score = min(score, max_score)
        print(f"REWARD: {final_score}")
        return final_score

    # Requirement 5: Comment positioned next to chart (0.2)
    chart_left, chart_top = chart.left, chart.top
    chart_width, chart_height = chart.width, chart.height
    comment_left, comment_top = comment.left, comment.top

    # Check horizontal proximity: comment to the right of chart (allow minor overlap)
    horizontal_ok = comment_left >= (chart_left + chart_width - int(chart_width * 0.1))
    # Check vertical alignment overlap
    vertical_ok = (comment_top + comment.height >= chart_top) and (comment_top <= chart_top + chart_height)

    if horizontal_ok and vertical_ok:
        print("✓ Comment placed next to chart (0.2)")
        score += 0.2
    else:
        print("✗ Comment not positioned next to chart")
        print(f"Positions -> chart: left={chart_left}, width={chart_width}, top={chart_top}, height={chart_height}")
        print(f"            comment: left={comment_left}, top={comment_top}")

    # Final score calculation
    final_score = min(score, max_score)
    print(f"Final score: {final_score}")
    return final_score

if __name__ == "__main__":
    file_path = "/home/user/on_slide_7_theres_a_chart_and_id_like_to_make_sure_whoever_reviews_it_checks_the_numbers_carefully_h.pptx"
    reward = verify_task(file_path)
    print(f"REWARD: {reward}")
