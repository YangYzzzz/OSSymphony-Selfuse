"""
FINAL REWARD SCRIPT - SUCCESS
Task: I’m tweaking the slide labeled “Features” in my LibreOffice Impress file. Could you walk me through adding a table with exactly 5 columns and 2 rows and making sure it’s centered perfectly on that slide?
Generated: 2025-09-10 13:41:35
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation


def _is_centered(shape, slide_width, slide_height, tolerance=50_000):
    """Return True if *shape* is centered within *tolerance* EMUs."""
    horiz_center = shape.left + shape.width / 2
    vert_center = shape.top + shape.height / 2
    dx = abs(horiz_center - slide_width / 2)
    dy = abs(vert_center - slide_height / 2)
    print(f"      Center offset – dx: {dx}, dy: {dy} (tolerance {tolerance})")
    return dx <= tolerance and dy <= tolerance


def verify_impress_table_task(file_path: str) -> float:
    """Verify the task requirements for the ‘Features’ slide.

    Requirements:
    1. A slide whose text includes the word “Features”.     (0.3 pts)
    2. On that slide a table with exactly 2 rows & 5 cols.  (0.4 pts)
    3. That table is perfectly centred on the slide.        (0.3 pts)

    Progressive scoring is applied; maximum total is 1.0.
    """
    max_score = 1.0
    score = 0.0
    print(f"Checking presentation: {file_path}\n")

    # ---------- prerequisite: load presentation (no points) ----------
    if not os.path.exists(file_path):
        print("✗ File not found.")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as exc:
        print(f"✗ Error opening presentation: {exc}")
        print("REWARD: 0.0")
        return 0.0

    # ---------- 1) locate the ‘Features’ slide ----------
    features_slide = None
    for idx, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = shape.text.strip()
                if txt:
                    slide_text.append(txt)
        if any("features" in t.lower() for t in slide_text):
            features_slide = slide
            print(f"✓ Found ‘Features’ slide at index {idx}")
            score += 0.3
            break

    if features_slide is None:
        print("✗ No slide containing the word ‘Features’ found.")
        print(f"REWARD: {score}")
        return score

    # ---------- 2) verify a 2×5 table exists on that slide ----------
    table_shape = None
    for shape in features_slide.shapes:
        # shape_type 19 corresponds to MSO_SHAPE_TYPE.TABLE
        if shape.shape_type == 19:
            tbl = shape.table
            rows, cols = len(tbl.rows), len(tbl.columns)
            print(f"  Found table – {rows} rows × {cols} columns")
            if rows == 2 and cols == 5:
                print("  ✓ Table dimensions are correct (2×5)")
                score += 0.4
                table_shape = shape
                break
            else:
                print("  ✗ Table dimensions incorrect – requirement not met yet")

    if table_shape is None:
        print("✗ Correctly-sized table not found on ‘Features’ slide.")
        print(f"REWARD: {score}")
        return score

    # ---------- 3) verify the table is centred ----------
    slide_w, slide_h = prs.slide_width, prs.slide_height
    if _is_centered(table_shape, slide_w, slide_h):
        print("  ✓ Table is centred on the slide")
        score += 0.3
    else:
        print("  ✗ Table is not centred on the slide")

    final_score = min(score, max_score)
    print(f"\nTotal score: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score


# ------------------ run verification when executed ------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/im_tweaking_the_slide_labeled_features_in_my_libreoffice_impress_file_could_you_walk_me_through_addi_golden.pptx"
    verify_impress_table_task(FILE_PATH)

