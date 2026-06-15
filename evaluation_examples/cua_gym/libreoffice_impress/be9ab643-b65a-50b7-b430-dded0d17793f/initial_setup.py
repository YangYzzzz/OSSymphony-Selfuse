"""
Initial Setup: Animate the horizontal bar chart on slide 3 with Wipe animation
Task ID: impress_ma_089
Domain: libreoffice_impress

Creates a 6-slide Performance Review presentation. Slide 3 has a horizontal
bar chart with 6 bars showing employee ratings. No animations are applied.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_089'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ============================
    # Slide 1: Title Slide
    # ============================
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Annual Performance Review 2025"
    slide1.placeholders[1].text = "Quarterly Assessment — Engineering Division"

    # ============================
    # Slide 2: Overview
    # ============================
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Review Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "This review covers the performance ratings for six key team members in the Engineering Division."
    p2a = body2.add_paragraph()
    p2a.text = "Ratings are on a scale of 1 to 10 based on project delivery, code quality, collaboration, and leadership."
    p2a.level = 0
    p2b = body2.add_paragraph()
    p2b.text = "The assessment period spans January through December 2025."
    p2b.level = 0
    p2c = body2.add_paragraph()
    p2c.text = "Each employee was evaluated by their direct manager and two peer reviewers."
    p2c.level = 0

    # ============================
    # Slide 3: Horizontal Bar Chart (6 bars, employee ratings)
    # ============================
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title manually
    add_text_box(slide3, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "Employee Performance Ratings", font_size=28, bold=True,
                 color=RGBColor(0x1F, 0x38, 0x64), alignment=PP_ALIGN.LEFT)

    # Chart data
    chart_data = CategoryChartData()
    chart_data.categories = [
        'Sarah Chen', 'Marcus Johnson', 'Aisha Patel',
        'David Kim', 'Elena Rodriguez', 'James Thompson'
    ]
    chart_data.add_series('Rating', (8.7, 9.2, 7.8, 8.4, 9.5, 7.1))

    # Add horizontal bar chart
    chart_shape = slide3.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(1.0), Inches(1.3), Inches(11.0), Inches(5.5),
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # Style the chart
    plot = chart.plots[0]
    plot.gap_width = 80
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)

    # Value axis (horizontal) range
    value_axis = chart.value_axis
    value_axis.minimum_scale = 0
    value_axis.maximum_scale = 10

    # ============================
    # Slide 4: Key Findings
    # ============================
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Key Findings"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Elena Rodriguez achieved the highest rating (9.5) with exceptional project delivery across three major launches."
    p4a = body4.add_paragraph()
    p4a.text = "Marcus Johnson showed significant improvement from last quarter, rising from 7.8 to 9.2."
    p4a.level = 0
    p4b = body4.add_paragraph()
    p4b.text = "James Thompson's rating of 7.1 reflects challenges in the Q3 migration project timeline."
    p4b.level = 0
    p4c = body4.add_paragraph()
    p4c.text = "Overall team average: 8.45, exceeding the division target of 8.0."
    p4c.level = 0

    # ============================
    # Slide 5: Action Items
    # ============================
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Action Items & Next Steps"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Schedule 1:1 development plan sessions with each team member by January 15, 2026."
    p5a = body5.add_paragraph()
    p5a.text = "Assign James Thompson a senior mentor for the Q1 infrastructure project."
    p5a.level = 0
    p5b = body5.add_paragraph()
    p5b.text = "Nominate Elena Rodriguez and Marcus Johnson for the Leadership Development Program."
    p5b.level = 0
    p5c = body5.add_paragraph()
    p5c.text = "Review and update performance metrics framework for 2026 cycle."
    p5c.level = 0

    # ============================
    # Slide 6: Thank You / Closing
    # ============================
    slide6 = prs.slides.add_slide(prs.slide_layouts[0])
    slide6.shapes.title.text = "Thank You"
    slide6.placeholders[1].text = "Questions & Discussion"

    # Save the presentation
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
