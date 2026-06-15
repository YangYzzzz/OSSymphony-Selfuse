"""
FINAL REWARD SCRIPT - SUCCESS
Task: Mark paragraph 11 with Language = None to skip spell checking.
Generated: 2025-10-17 07:56:00
Status: success
Model: azure-o3
Total Steps: 12
"""

import os
import re
import zipfile
from pptx import Presentation
from pptx.enum.lang import MSO_LANGUAGE_ID

def verify_task(file_path: str) -> float:
    """Verify that paragraph 11 in the presentation is marked to skip spell checking.

    Scoring (progressive):
        0.3 – Target paragraph (text == "Paragraph 11") is found
        0.7 – Paragraph is actually marked to skip spell-checking via either
              a:noProof="1" attribute OR language_id == MSO_LANGUAGE_ID.NONE
        1.0 – Both conditions satisfied
    """

    total_score = 0.0
    max_score   = 1.0

    print(f"Verifying file: {file_path}")

    # 1. Basic existence & load (NO POINTS, but prerequisite)
    if not os.path.exists(file_path):
        print("✗ File not found")
        return 0.0
    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded successfully with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Could not load presentation: {e}")
        return 0.0

    # 2. Locate the paragraph whose visible text is exactly "Paragraph 11"
    target_para = None
    para_counter = 0  # overall paragraph index for information only

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                para_counter += 1
                text = "".join(run.text for run in para.runs).strip()
                if text.lower() == "paragraph 11":
                    target_para = para
                    target_index = para_counter  # 1-based index for printing
                    break
            if target_para:
                break
        if target_para:
            break

    if target_para is None:
        print("✗ Required paragraph with text 'Paragraph 11' not found")
        return 0.0

    print(f"✓ Found target paragraph at overall index {target_index}")
    total_score += 0.3  # paragraph located

    # 3. Verify spell-check suppression markers
    lang_none_found = False
    no_proof_found  = False

    for run in target_para.runs:
        # a) Language ID check via python-pptx API
        if run.font.language_id == MSO_LANGUAGE_ID.NONE:
            lang_none_found = True
        # b) Raw XML check for a:noProof="1"
        if "noProof=\"1\"" in run._r.xml or "noProof=\"true\"" in run._r.xml:
            no_proof_found = True

    if lang_none_found or no_proof_found:
        print("✓ Target paragraph is marked to skip spell checking")
        total_score += 0.7  # full additional credit
    else:
        print("✗ Target paragraph is NOT marked to skip spell checking")

    final_score = min(total_score, max_score)
    print(f"REWARD: {final_score}")
    return final_score

# ---------------------------------------------------------------------------
# Execute verification when script is run directly
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/mark_paragraph_11_with_language_none_to_skip_spell_checking.pptx"
    verify_task(FILE_PATH)
