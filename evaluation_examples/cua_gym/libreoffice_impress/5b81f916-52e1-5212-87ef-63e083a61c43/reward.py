"""
Reward Script: Insert right-pointing block arrow on slide 6 with gradient and formatted text
Task ID: impress_ndo_055
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Right arrow shape exists on slide 6
  Component 2 (0.20): Size approximately 8cm x 3cm
  Component 3 (0.20): Gradient fill #2E86C1 -> #1B4F72 left-to-right
  Component 4 (0.20): Text is 'Next Step'
  Component 5 (0.15): Text white, 16pt, bold, centered
"""

import os
from pptx import Presentation
from pptx.util import Emu, Pt, Cm
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_055'


def is_approx_equal(val1, val2, tolerance=0.05):
    """Check approximate equality with relative tolerance."""
    if val1 == val2:
        return True
    if val1 == 0 or val2 == 0:
        return abs(val1 - val2) < 10000  # small absolute tolerance for zero
    return abs(val1 - val2) / max(abs(val1), abs(val2)) <= tolerance


def find_right_arrow_on_slide6(prs):
    """Find a right arrow shape on slide 6 (index 5)."""
    if len(prs.slides) < 6:
        return None
    slide = prs.slides[5]
    for shape in slide.shapes:
        prstGeom = shape._element.find('.//' + qn('a:prstGeom'))
        if prstGeom is not None:
            prst = prstGeom.get('prst')
            if prst == 'rightArrow':
                return shape
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Right arrow shape exists on slide 6 (0.25 points)
    arrow_shape = None
    try:
        arrow_shape = find_right_arrow_on_slide6(prs)
        if arrow_shape is not None:
            print(f"PASS: Component 1 -- Right arrow found on slide 6 (name={arrow_shape.name}) (0.25 pts)")
            total_score += 0.25
        else:
            print("FAIL: Component 1 -- No right arrow shape found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    if arrow_shape is None:
        # No arrow found, remaining checks are meaningless
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: Size approximately 8cm x 3cm (0.20 points)
    try:
        expected_w = Cm(8)   # 2880000 EMU
        expected_h = Cm(3)   # 1080000 EMU
        actual_w = arrow_shape.width
        actual_h = arrow_shape.height
        w_ok = is_approx_equal(actual_w, expected_w)
        h_ok = is_approx_equal(actual_h, expected_h)
        if w_ok and h_ok:
            print(f"PASS: Component 2 -- Size {actual_w/360000:.2f}cm x {actual_h/360000:.2f}cm matches expected 8cm x 3cm (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 2 -- Size {actual_w/360000:.2f}cm x {actual_h/360000:.2f}cm, expected ~8cm x ~3cm (w_ok={w_ok}, h_ok={h_ok})")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Gradient fill from #2E86C1 to #1B4F72 left-to-right (0.20 points)
    try:
        elem = arrow_shape._element
        gradFill = elem.find('.//' + qn('a:gradFill'))
        if gradFill is not None:
            # Extract gradient stops
            gs_list = gradFill.findall('.//' + qn('a:gs'))
            colors = {}
            for gs in gs_list:
                pos = int(gs.get('pos', '0'))
                srgbClr = gs.find(qn('a:srgbClr'))
                if srgbClr is not None:
                    colors[pos] = srgbClr.get('val', '').upper()

            # Check linear direction (ang=0 means left-to-right)
            lin = gradFill.find(qn('a:lin'))
            ang = None
            if lin is not None:
                ang = int(lin.get('ang', '-1'))

            # Verify: start (pos 0) = 2E86C1, end (pos 100000) = 1B4F72, linear angle 0
            start_color = colors.get(0, '').upper()
            end_color = colors.get(100000, '').upper()
            color_ok = (start_color == '2E86C1' and end_color == '1B4F72')
            direction_ok = (ang == 0)

            if color_ok and direction_ok:
                print(f"PASS: Component 3 -- Gradient fill #2E86C1 -> #1B4F72, angle=0 (left-to-right) (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 3 -- start={start_color}, end={end_color}, angle={ang} (expected 2E86C1->1B4F72, angle=0)")
        else:
            print("FAIL: Component 3 -- No gradient fill found on arrow shape")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Text is 'Next Step' (0.20 points)
    try:
        if hasattr(arrow_shape, 'text_frame'):
            full_text = arrow_shape.text_frame.text.strip()
            if full_text == 'Next Step':
                print(f"PASS: Component 4 -- Text is 'Next Step' (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 4 -- Text is {repr(full_text)}, expected 'Next Step'")
        else:
            print("FAIL: Component 4 -- Arrow shape has no text frame")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Text formatting: white (#FFFFFF), 16pt, bold, centered (0.15 points)
    try:
        tf = arrow_shape.text_frame
        # Find the run with 'Next Step' text
        target_run = None
        target_para = None
        for para in tf.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    target_run = run
                    target_para = para
                    break
            if target_run:
                break

        if target_run is None:
            print("FAIL: Component 5 -- No text run found in arrow")
        else:
            checks_passed = 0
            total_checks = 4

            # Check bold
            if target_run.font.bold:
                checks_passed += 1
            else:
                print(f"  DETAIL: bold={target_run.font.bold}, expected True")

            # Check size ~16pt (203200 EMU)
            expected_size = Pt(16)  # 203200 EMU
            actual_size = target_run.font.size
            if actual_size is not None and is_approx_equal(actual_size, expected_size, tolerance=0.05):
                checks_passed += 1
            else:
                print(f"  DETAIL: size={actual_size}, expected {expected_size}")

            # Check white color
            try:
                rgb = target_run.font.color.rgb
                if str(rgb).upper() == 'FFFFFF':
                    checks_passed += 1
                else:
                    print(f"  DETAIL: color={rgb}, expected FFFFFF")
            except Exception:
                print(f"  DETAIL: could not read font color")

            # Check centered alignment
            from pptx.enum.text import PP_ALIGN
            align = target_para.alignment
            if align == PP_ALIGN.CENTER:
                checks_passed += 1
            else:
                print(f"  DETAIL: alignment={align}, expected CENTER")

            if checks_passed == total_checks:
                print(f"PASS: Component 5 -- Text is white, 16pt, bold, centered (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 5 -- {checks_passed}/{total_checks} formatting checks passed")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
