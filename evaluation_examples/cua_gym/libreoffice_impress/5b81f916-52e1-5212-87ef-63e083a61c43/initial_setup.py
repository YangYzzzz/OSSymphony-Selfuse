"""
Initial Setup: Process presentation with 8 slides; slide 6 has upper content, empty bottom area.
Task ID: impress_ndo_055
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_055'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
    return slide


def add_blank_with_title(prs, title_text):
    """Add a slide with title only layout (layout index 5=blank)."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Cm(1.5), Cm(0.8), Cm(22), Cm(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1B, 0x4F, 0x72)
    return slide


def create_initial():
    prs = Presentation()

    # Slide 1: Title slide
    add_title_slide(prs, "Project Workflow Process",
                    "Q2 2025 Operations Review — Riverside Analytics")

    # Slide 2: Overview
    add_content_slide(prs, "Process Overview", [
        "Our workflow consists of seven key phases",
        "Each phase has defined entry and exit criteria",
        "Stakeholder review gates between major phases",
        "Total cycle time target: 12 business days",
    ])

    # Slide 3: Phase 1 & 2
    add_content_slide(prs, "Requirements & Planning", [
        "Phase 1 — Requirements Gathering (2 days)",
        "  Conduct stakeholder interviews with product team",
        "  Document functional and non-functional requirements",
        "Phase 2 — Sprint Planning (1 day)",
        "  Prioritize backlog items by business value",
        "  Assign story points and team capacity",
    ])

    # Slide 4: Phase 3 & 4
    add_content_slide(prs, "Design & Development", [
        "Phase 3 — Architecture Design (2 days)",
        "  Create system design documents",
        "  Review with senior engineering leads",
        "Phase 4 — Implementation (4 days)",
        "  Follow coding standards and branch strategy",
        "  Daily stand-ups and pair programming sessions",
    ])

    # Slide 5: Phase 5
    add_content_slide(prs, "Quality Assurance", [
        "Phase 5 — Testing & Validation (2 days)",
        "  Unit tests: minimum 85% code coverage",
        "  Integration tests across all service boundaries",
        "  Performance benchmarks against SLA targets",
        "  Security scan with OWASP top-10 checklist",
    ])

    # Slide 6: Phase 6 — Deployment (upper portion has content, bottom empty)
    # Using blank layout so we control placement in the upper area only
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])

    # Title in upper area
    title_box = slide6.shapes.add_textbox(Cm(1.5), Cm(0.8), Cm(22), Cm(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Deployment & Release"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1B, 0x4F, 0x72)

    # Content in upper-middle area
    content_box = slide6.shapes.add_textbox(Cm(2), Cm(3.5), Cm(20), Cm(7))
    tf2 = content_box.text_frame
    tf2.word_wrap = True

    lines = [
        "Phase 6 — Deployment Pipeline",
        "  Blue-green deployment to staging environment",
        "  Automated smoke tests run within 5 minutes",
        "  Manual UAT sign-off from product owner",
        "  Rolling release to production clusters",
        "  Monitoring dashboards active for 24h post-deploy",
    ]
    for i, line in enumerate(lines):
        if i == 0:
            tf2.paragraphs[0].text = line
            r = tf2.paragraphs[0].runs[0]
            r.font.size = Pt(18)
            r.font.bold = True
        else:
            para = tf2.add_paragraph()
            para.text = line
            r = para.runs[0]
            r.font.size = Pt(14)

    # NOTE: Bottom area of slide 6 is intentionally left empty
    # The task will ask to add a block arrow here

    # Slide 7: Phase 7
    add_content_slide(prs, "Post-Release Monitoring", [
        "Phase 7 — Observation & Feedback (1 day)",
        "  Monitor error rates and latency metrics",
        "  Collect user feedback through support channels",
        "  Run retrospective meeting with full team",
        "  Document lessons learned in Confluence wiki",
    ])

    # Slide 8: Summary
    add_content_slide(prs, "Process Summary & Next Steps", [
        "Total estimated cycle: 12 business days",
        "Key milestones tracked in Jira dashboard",
        "Quarterly process improvement reviews scheduled",
        "Contact: ops-team@riverside-analytics.com",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
