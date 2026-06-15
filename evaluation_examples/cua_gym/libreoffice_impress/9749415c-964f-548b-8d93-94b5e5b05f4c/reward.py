"""
FINAL REWARD SCRIPT - SUCCESS
Task: While polishing my deck in LibreOffice Impress, I want to flag a slide for the Q&A section: how do I add the speaker note "APP" to one slide and switch that same slide’s background fill to solid purple (#800080) so it stands out at a glance?
Generated: 2025-09-10 13:11:23
Status: success
Model: azure-o3
Total Steps: 1
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
import os


def _is_purple(rgb):
    """Return True only if the given RGBColor is exactly #800080."""
    return rgb == RGBColor(0x80, 0x00, 0x80)


def _slide_has_app_note(slide):
    """Check if the slide’s notes contain the word APP (case-insensitive)."""
    if not slide.has_notes_slide:
        return False
    try:
        notes_slide = slide.notes_slide
        texts = []
        for shape in notes_slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        notes_text = " ".join(texts).lower()
        return "app" in notes_text
    except Exception:
        return False


def _slide_has_purple_bg(slide):
    """Check if the slide background is set to a SOLID fill of #800080."""
    try:
        fill = slide.background.fill
        if fill and fill.type == MSO_FILL.SOLID:
            return _is_purple(fill.fore_color.rgb)
    except Exception:
        pass
    return False


def verify_impress_task(file_path):
    """Verify the task requirements and return a progressive score (0.0-1.0)."""
    print(f"Verifying presentation: {file_path}")

    # 0. PRECONDITION: file must exist and load
    if not os.path.exists(file_path):
        print("✗ File not found")
        print("REWARD: 0.0")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Could not load PPTX ({e})")
        print("REWARD: 0.0")
        return 0.0

    # 1. ANALYSE every slide for the two requirements
    slides_with_purple = []
    slides_with_app    = []
    slides_with_both   = []

    for idx, slide in enumerate(prs.slides, 1):
        has_purple = _slide_has_purple_bg(slide)
        has_app    = _slide_has_app_note(slide)

        if has_purple:
            slides_with_purple.append(idx)
            print(f"  ✓ Slide {idx}: solid purple background detected")
        if has_app:
            slides_with_app.append(idx)
            print(f"  ✓ Slide {idx}: speaker notes contain 'APP'")
        if has_purple and has_app:
            slides_with_both.append(idx)

    # 2. SCORING (progressive)
    score = 0.0
    if slides_with_both:
        print(f"✓ Slide(s) {slides_with_both} satisfy BOTH requirements – perfect.")
        score = 1.0
    else:
        # partial credit if only one aspect met (each worth 0.5)
        if slides_with_purple:
            score += 0.5
        if slides_with_app:
            score += 0.5

    score = min(score, 1.0)
    print(f"REWARD: {score}")
    return score


# ---------------------------------------------------------------------------
# When the script is executed on the VM, run the verification automatically
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    FILE = "/home/user/while_polishing_my_deck_in_libreoffice_impress_i_want_to_flag_a_slide_for_the_qa_section_how_do_i_ad_golden.pptx"
    verify_impress_task(FILE)

