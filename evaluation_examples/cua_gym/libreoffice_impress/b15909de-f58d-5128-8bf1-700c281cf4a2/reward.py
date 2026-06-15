"""
FINAL REWARD SCRIPT - SUCCESS
Task: I’m putting the finishing touches on my deck and just noticed slide 105 doesn’t have any speaker notes. In LibreOffice Impress, how do I bring up the Notes view for that slide and type the line “Emphasize results”?
Generated: 2025-09-10 23:38:29
Status: success
Model: azure-o3
Total Steps: 1
"""

from pptx import Presentation
import os, traceback

def verify_notes_on_slide(file_path: str,
                          slide_number: int = 105,
                          expected_phrase: str = "Emphasize results") -> float:
    """Verify that a specific slide in a PPTX has a notes slide
    containing a given phrase.

    Parameters
    ----------
    file_path : str
        Full path to the presentation file.
    slide_number : int, optional
        1-based index of the slide to inspect (default is 105).
    expected_phrase : str, optional
        Phrase that must appear in the speaker notes (case-insensitive).

    Returns
    -------
    float
        Progressive score between 0.0 and 1.0.
    """

    print(f"Verifying notes for slide {slide_number} in file: {file_path}\n")

    total_score = 0.0        # progressive score
    max_score   = 1.0        # cap

    # ----------------------------------------------------
    # 1. File existence & loading (no points, prerequisite)
    # ----------------------------------------------------
    if not os.path.exists(file_path):
        print(f"✗ File does not exist: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded successfully with {len(prs.slides)} slides")
    except Exception:
        print("✗ Failed to load presentation:\n", traceback.format_exc())
        print("REWARD: 0.0")
        return 0.0

    # --------------------------------------------
    # 2. Ensure requested slide exists (no points)
    # --------------------------------------------
    if len(prs.slides) < slide_number:
        print(f"✗ Presentation has only {len(prs.slides)} slides (< {slide_number}).")
        print("REWARD: 0.0")
        return 0.0

    target_slide = prs.slides[slide_number - 1]

    # -------------------------------------------
    # 3. Verify presence of a notes slide (0.5 pt)
    # -------------------------------------------
    notes_text_list = []
    if target_slide.has_notes_slide:
        notes_slide = target_slide.notes_slide
        # extract all text from the notes slide
        for shape in notes_slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text:
                        cleaned = paragraph.text.strip()
                        if cleaned:
                            notes_text_list.append(cleaned)
        print("✓ Notes slide exists for the target slide")
        total_score += 0.5
    else:
        print("✗ No notes slide found for the target slide")

    # -----------------------------------------------------
    # 4. Verify expected phrase inside the notes (0.5 pts)
    # -----------------------------------------------------
    expected_lower = expected_phrase.lower()
    if any(expected_lower in text.lower() for text in notes_text_list):
        print(f"✓ Found expected phrase '{expected_phrase}' in notes")
        total_score += 0.5
    else:
        print(f"✗ Expected phrase '{expected_phrase}' NOT found in notes")

    # --------------------
    # 5. Finalise and exit
    # --------------------
    final_score = min(total_score, max_score)
    print(f"\nTotal score: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score

# ----------------------
# Main script execution
# ----------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/im_putting_the_finishing_touches_on_my_deck_and_just_noticed_slide_105_doesnt_have_any_speaker_notes_golden.pptx"
    verify_notes_on_slide(FILE_PATH)

