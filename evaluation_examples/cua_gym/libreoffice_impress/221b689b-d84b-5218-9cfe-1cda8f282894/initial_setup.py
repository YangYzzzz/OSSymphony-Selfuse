"""
Initial Setup: Create a 24-slide all-hands presentation
Task ID: impress_tm_094
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_094'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points, layout_idx=1):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title_text
    body = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            body = ph
            break
    if body and body.has_text_frame:
        tf = body.text_frame
        tf.clear()
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.level = 0
            for run in p.runs:
                run.font.size = Pt(18)
    return slide


def set_slide_bg(slide, r, g, b):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    s = add_title_slide(prs, "Q1 2025 All-Hands Meeting", "Acme Technologies Inc.\nMarch 28, 2025")
    set_slide_bg(s, 0x1B, 0x3A, 0x5C)
    for shape in s.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Slide 2: Agenda
    add_content_slide(prs, "Agenda", [
        "Company Performance Overview",
        "Product Roadmap Update",
        "Engineering Milestones",
        "Sales & Revenue Highlights",
        "Customer Success Stories",
        "HR & Culture Initiatives",
        "Q&A Session"
    ])

    # Slide 3: Company Performance
    add_content_slide(prs, "Company Performance", [
        "Revenue: $47.2M (+23% YoY)",
        "Active users: 1.2M (+340K new)",
        "Net retention rate: 118%",
        "Employee count: 485 (+62 in Q1)"
    ])

    # Slide 4: Revenue Breakdown
    add_content_slide(prs, "Revenue Breakdown by Segment", [
        "Enterprise: $28.3M (60%)",
        "Mid-Market: $12.1M (26%)",
        "SMB: $4.7M (10%)",
        "Government: $2.1M (4%)"
    ])

    # Slide 5: Product Roadmap
    add_content_slide(prs, "Product Roadmap - H1 2025", [
        "AI-Powered Analytics Dashboard (Launch: April 15)",
        "Mobile App v3.0 with offline mode (Beta: May 1)",
        "Enterprise SSO Integration (GA: March 30)",
        "Real-time Collaboration Features (Q2)"
    ])

    # Slide 6: Engineering
    add_content_slide(prs, "Engineering Milestones", [
        "Migrated 92% of services to Kubernetes",
        "Reduced API latency p99 from 340ms to 89ms",
        "Achieved 99.97% uptime in Q1",
        "Deployed 847 PRs across 12 repositories"
    ])

    # Slide 7: Architecture
    add_content_slide(prs, "Architecture Modernization", [
        "Microservices migration: 18 of 23 services complete",
        "Database sharding for Analytics service",
        "CDN optimization: 40% reduction in load times",
        "New CI/CD pipeline: 3x faster deployments"
    ])

    # Slide 8: Sales Highlights
    add_content_slide(prs, "Sales & Revenue Highlights", [
        "52 new enterprise deals closed in Q1",
        "Average deal size increased to $142K (+18%)",
        "Pipeline value: $89M (2.1x coverage ratio)",
        "Win rate improved to 34% from 28%"
    ])

    # Slide 9: Top Deals
    add_content_slide(prs, "Notable Q1 Deals", [
        "GlobalBank Corp - $2.4M (3-year, 15K seats)",
        "Pacific Health Systems - $1.8M (enterprise suite)",
        "NovaTech Industries - $950K (analytics platform)",
        "Meridian Logistics - $720K (expansion deal)"
    ])

    # Slide 10: Customer Success
    add_content_slide(prs, "Customer Success Metrics", [
        "NPS Score: 72 (up from 65 in Q4)",
        "Support ticket resolution: 4.2 hrs avg (-31%)",
        "Customer health score: 8.4/10 average",
        "Churn rate: 1.2% (down from 1.8%)"
    ])

    # Slide 11: Case Study
    add_content_slide(prs, "Case Study: RetailMax Implementation", [
        "Challenge: Manual inventory across 200+ stores",
        "Solution: Real-time analytics + automated alerts",
        "Results: 34% reduction in stockouts",
        "ROI: 280% in first 6 months"
    ])

    # Slide 12: Marketing
    add_content_slide(prs, "Marketing Highlights", [
        "Website traffic: 2.8M visits (+45% QoQ)",
        "Content downloads: 12,400 whitepapers",
        "Webinar attendance: 3,200 registrants",
        "Brand awareness up 22% in target segments"
    ])

    # Slide 13: Events
    add_content_slide(prs, "Upcoming Events & Conferences", [
        "SaaSTech Summit - San Francisco (April 8-10)",
        "European Digital Conference - Berlin (April 22)",
        "Customer Advisory Board - NYC (May 5)",
        "AcmeCon 2025 - Austin (June 15-17)"
    ])

    # Slide 14: HR
    add_content_slide(prs, "People & Culture Update", [
        "62 new hires in Q1 (Engineering: 28, Sales: 18)",
        "Employee satisfaction score: 4.3/5.0",
        "Voluntary turnover: 8.2% annualized",
        "DEI initiatives: 6 new ERGs launched"
    ])

    # Slide 15: Learning
    add_content_slide(prs, "Learning & Development", [
        "$2,500 annual learning stipend (new!)",
        "Internal mentorship program: 120 pairs matched",
        "Leadership development cohort: 24 participants",
        "Technical certifications funded: 89 in Q1"
    ])

    # Slide 16: Benefits
    add_content_slide(prs, "New Benefits & Perks", [
        "Extended parental leave: 16 weeks paid",
        "Wellness reimbursement increased to $1,200/yr",
        "Hybrid work policy: 3 days flexible",
        "Sabbatical program after 5 years of service"
    ])

    # Slide 17: Finance
    add_content_slide(prs, "Financial Outlook - Q2 2025", [
        "Revenue target: $52M (+10% QoQ)",
        "Planned headcount additions: 45",
        "R&D investment increase: 15%",
        "Expected EBITDA margin: 12%"
    ])

    # Slide 18: Ops
    add_content_slide(prs, "Operations & Infrastructure", [
        "New data center in Frankfurt (live April 1)",
        "SOC 2 Type II audit completed successfully",
        "ISO 27001 certification in progress",
        "Disaster recovery RTO reduced to 15 minutes"
    ])

    # Slide 19: Security
    add_content_slide(prs, "Security & Compliance", [
        "Zero critical security incidents in Q1",
        "Penetration testing: all findings remediated",
        "GDPR data processing agreements: 100% coverage",
        "Security awareness training: 98% completion"
    ])

    # Slide 20: Partnerships
    add_content_slide(prs, "Strategic Partnerships", [
        "AWS Advanced Technology Partner status achieved",
        "Salesforce AppExchange integration launched",
        "Microsoft Teams connector in beta",
        "New channel partner program: 14 partners signed"
    ])

    # Slide 21: Innovation
    add_content_slide(prs, "Innovation Lab Projects", [
        "Project Aurora: Generative AI for report creation",
        "Project Nexus: Predictive customer health scoring",
        "Project Titan: Edge computing for IoT data",
        "Hackathon Q1: 32 projects, 3 moved to roadmap"
    ])

    # Slide 22: Community
    add_content_slide(prs, "Community & Social Impact", [
        "Volunteer hours logged: 1,240 in Q1",
        "STEM education partnership with 8 local schools",
        "Carbon offset program: 120% of emissions",
        "$50K donated to tech education nonprofits"
    ])

    # Slide 23: Key Priorities
    s23 = add_content_slide(prs, "Key Priorities for Q2 2025", [
        "Ship AI Analytics Dashboard by April 15",
        "Close 60+ enterprise deals",
        "Achieve $52M quarterly revenue target",
        "Complete Kubernetes migration (100%)",
        "Launch mobile app v3.0 beta"
    ])

    # Slide 24: Closing
    s24 = add_title_slide(prs, "Thank You!", "Questions & Discussion\nSlack: #all-hands-q1-2025")
    set_slide_bg(s24, 0x1B, 0x3A, 0x5C)
    for shape in s24.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # Ensure Documents directory exists
    os.makedirs(f'{WORKDIR}/Documents', exist_ok=True)

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
