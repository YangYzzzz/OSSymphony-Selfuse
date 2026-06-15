"""
Initial Setup: Add a hyperlink in the notes of slide 5
Task ID: impress_ndo_018
Domain: libreoffice_impress

Creates a 6-slide research presentation. Slide 5 has speaker notes with
plain text only (no hyperlink). The agent must add the hyperlink.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_018'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Market Research Findings Q4 2025"
    slide1.placeholders[1].text = "Prepared by: Analytics Division\nDecember 2025"

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Our Q4 research covered three primary sectors:"
    for item in [
        "Consumer electronics saw a 12% increase in demand",
        "Cloud services grew by 23% year-over-year",
        "Healthcare technology adoption reached 67% penetration",
    ]:
        p = tf2.add_paragraph()
        p.text = item
        p.level = 1
    slide2.notes_slide.notes_text_frame.text = (
        "Emphasize the cloud services growth as the main headline."
    )

    # --- Slide 3: Methodology ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Research Methodology"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Data collection ran from October 1 to November 30, 2025."
    for item in [
        "Survey sample: 4,200 respondents across 8 regions",
        "Interviews: 35 industry leaders and CTO-level executives",
        "Public datasets: SEC filings, Gartner reports, IDC forecasts",
        "Confidence interval: 95% with a margin of error of +/- 2.3%",
    ]:
        p = tf3.add_paragraph()
        p.text = item
        p.level = 1
    slide3.notes_slide.notes_text_frame.text = (
        "Mention that the sample size was increased by 15% compared to Q3."
    )

    # --- Slide 4: Key Findings ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Key Findings"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Top insights from the research:"
    findings = [
        "Mobile-first strategies dominate 78% of new product launches",
        "AI-powered analytics tools reduced decision latency by 34%",
        "Customer retention improved 19% with personalized outreach",
        "Supply chain automation cut costs by $2.1M on average",
        "Remote workforce tools saw 41% higher engagement scores",
    ]
    for item in findings:
        p = tf4.add_paragraph()
        p.text = item
        p.level = 1
    slide4.notes_slide.notes_text_frame.text = (
        "Pause after the AI analytics finding to allow questions."
    )

    # --- Slide 5: Data Sources & References ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Data Sources & References"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Primary and secondary sources used in this report:"
    sources = [
        "Gartner Market Analysis Report, November 2025",
        "IDC Worldwide Spending Guide, Q4 2025 Edition",
        "Internal CRM data, October-November 2025",
        "Bureau of Labor Statistics employment data",
        "McKinsey Global Institute technology adoption survey",
    ]
    for item in sources:
        p = tf5.add_paragraph()
        p.text = item
        p.level = 1

    # Slide 5 notes: plain text only, NO hyperlink
    slide5.notes_slide.notes_text_frame.text = (
        "Cite the following source during this section."
    )

    # --- Slide 6: Next Steps ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Next Steps"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Recommended actions for Q1 2026:"
    for item in [
        "Expand survey coverage to include APAC markets",
        "Commission a deep-dive study on AI analytics ROI",
        "Schedule stakeholder review meeting by January 15",
    ]:
        p = tf6.add_paragraph()
        p.text = item
        p.level = 1
    slide6.notes_slide.notes_text_frame.text = (
        "Close with a call to action for budget approval."
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
