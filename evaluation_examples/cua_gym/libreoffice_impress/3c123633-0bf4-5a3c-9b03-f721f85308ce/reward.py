"""
Reward Script: Add Fade Out exit animations to three images on slide 4
Task ID: impress_ma_061
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): Animation timing exists with 3 exit animations on slide 4
  Component 2 (0.3): All three animations target correct shapes with Fade exit effect
  Component 3 (0.2): Correct triggers — Photo1=clickEffect, Photo2/3=afterEffect
  Component 4 (0.2): Correct 0.5s delay between sequential animations
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_061'

NS = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
}


def get_shape_name_to_id(pptx_path, slide_idx):
    """Get shape name to spid mapping from slide XML."""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_idx]
    mapping = {}
    for shape in slide.shapes:
        mapping[shape.name] = shape.shape_id
    return mapping


def parse_animations(pptx_path, slide_idx):
    """Parse animation data from slide XML. Returns list of animation dicts."""
    animations = []
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        slide_name = f'ppt/slides/slide{slide_idx + 1}.xml'
        try:
            with zf.open(slide_name) as f:
                root = ET.fromstring(f.read())
        except KeyError:
            return animations

    timing = root.find('.//p:timing', NS)
    if timing is None:
        return animations

    # Find all cTn elements with presetClass attribute (these are the animation effects)
    for ctn in timing.iter('{http://schemas.openxmlformats.org/presentationml/2006/main}cTn'):
        pass  # iter doesn't work well across namespaces

    # Use full namespace iteration
    pns = 'http://schemas.openxmlformats.org/presentationml/2006/main'

    # Find all par elements that contain animation presets
    for ctn_el in root.iter(f'{{{pns}}}cTn'):
        preset_id = ctn_el.get('presetID')
        preset_class = ctn_el.get('presetClass')
        node_type = ctn_el.get('nodeType')

        if preset_id and preset_class:
            # Find the target shape ID from child animEffect -> cBhvr -> tgtEl -> spTgt
            spid = None
            filter_val = None
            anim_dur = None

            for anim_effect in ctn_el.iter(f'{{{pns}}}animEffect'):
                filter_val = anim_effect.get('filter')
                trans = anim_effect.get('transition')
                for cbhvr in anim_effect.iter(f'{{{pns}}}cBhvr'):
                    for tgt_el in cbhvr.iter(f'{{{pns}}}tgtEl'):
                        for sp_tgt in tgt_el.iter(f'{{{pns}}}spTgt'):
                            spid = sp_tgt.get('spid')
                    for inner_ctn in cbhvr.iter(f'{{{pns}}}cTn'):
                        if inner_ctn.get('dur'):
                            anim_dur = inner_ctn.get('dur')

            # Get delay from parent's stCondLst
            parent_ctn = ctn_el
            delay = None
            st_cond_list = ctn_el.find(f'{{{pns}}}stCondLst')
            if st_cond_list is not None:
                for cond in st_cond_list.iter(f'{{{pns}}}cond'):
                    delay = cond.get('delay')

            animations.append({
                'preset_id': preset_id,
                'preset_class': preset_class,
                'node_type': node_type,
                'spid': spid,
                'filter': filter_val,
                'duration': anim_dur,
                'delay': delay,
            })

    return animations


def get_parent_delays(pptx_path, slide_idx):
    """Parse parent-level delay conditions for After Previous animations.
    Returns dict of spid -> delay_ms from parent container's stCondLst."""
    delays = {}
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        slide_name = f'ppt/slides/slide{slide_idx + 1}.xml'
        try:
            with zf.open(slide_name) as f:
                content = f.read().decode()
        except KeyError:
            return delays

    root = ET.fromstring(content)
    pns = 'http://schemas.openxmlformats.org/presentationml/2006/main'

    # Find all par elements that wrap afterEffect animations
    # The delay is on the grandparent par's cTn stCondLst
    for par_el in root.iter(f'{{{pns}}}par'):
        ctn = par_el.find(f'{{{pns}}}cTn')
        if ctn is None:
            continue

        # Check if this par contains an afterEffect animation
        target_spid = None
        found_after = any(
            inner_ctn.get('nodeType') == 'afterEffect'
            for inner_ctn in ctn.iter(f'{{{pns}}}cTn')
        )
        if found_after:
            for inner_ctn in ctn.iter(f'{{{pns}}}cTn'):
                if inner_ctn.get('nodeType') == 'afterEffect':
                    for sp_tgt in inner_ctn.iter(f'{{{pns}}}spTgt'):
                        target_spid = sp_tgt.get('spid')
                    break

        if not found_after or not target_spid:
            continue

        # The delay is on this par's cTn stCondLst (the wrapping container)
        st_cond = ctn.find(f'{{{pns}}}stCondLst')
        if st_cond is not None:
            for cond in st_cond.iter(f'{{{pns}}}cond'):
                d = cond.get('delay')
                evt = cond.get('evt')
                if d and evt == 'onEnd':
                    delays[target_spid] = int(d)

    return delays


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Get shape name -> id mapping
    try:
        name_to_id = get_shape_name_to_id(file_path, 3)  # slide 4 (0-indexed)
        print(f"Shape mapping: {name_to_id}")
    except Exception as e:
        print(f"CRITICAL: Cannot load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Expected shape IDs for Photo1, Photo2, Photo3
    photo1_id = name_to_id.get('Photo1')
    photo2_id = name_to_id.get('Photo2')
    photo3_id = name_to_id.get('Photo3')

    if not all([photo1_id, photo2_id, photo3_id]):
        print(f"CRITICAL: Missing photo shapes. Found: Photo1={photo1_id}, Photo2={photo2_id}, Photo3={photo3_id}")
        print("REWARD: 0.0")
        return 0.0

    expected_spids = {str(photo1_id), str(photo2_id), str(photo3_id)}

    # Parse animations
    try:
        animations = parse_animations(file_path, 3)
        print(f"Found {len(animations)} animation presets")
        for a in animations:
            print(f"  preset={a['preset_id']}, class={a['preset_class']}, "
                  f"node={a['node_type']}, spid={a['spid']}, filter={a['filter']}")
    except Exception as e:
        print(f"CRITICAL: Cannot parse animations: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Filter to exit animations only
    exit_anims = [a for a in animations if a['preset_class'] == 'exit']
    print(f"Exit animations: {len(exit_anims)}")

    # Component 1: Animation timing exists with 3 exit animations on slide 4 (0.3 points)
    try:
        if len(exit_anims) >= 3:
            print(f"PASS: Component 1 — Found {len(exit_anims)} exit animations on slide 4 (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 — Expected 3 exit animations, found {len(exit_anims)}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: All three animations target Photo1/2/3 with Fade effect (0.3 points)
    try:
        targeted_spids = set()
        fade_count = 0
        for a in exit_anims:
            if a['spid'] in expected_spids:
                targeted_spids.add(a['spid'])
                # presetID=10 is Fade, filter should be 'fade'
                if a['preset_id'] == '10' and a.get('filter') == 'fade':
                    fade_count += 1

        if targeted_spids == expected_spids and fade_count >= 3:
            print(f"PASS: Component 2 — All 3 photos have Fade exit effect (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 — Targeted shapes: {targeted_spids} (expected {expected_spids}), "
                  f"fade count: {fade_count}/3")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Correct triggers — Photo1=clickEffect/withEffect, Photo2/3=afterEffect (0.2 points)
    try:
        anim_by_spid = {}
        for a in exit_anims:
            if a['spid'] in expected_spids:
                anim_by_spid[a['spid']] = a

        photo1_anim = anim_by_spid.get(str(photo1_id), {})
        photo2_anim = anim_by_spid.get(str(photo2_id), {})
        photo3_anim = anim_by_spid.get(str(photo3_id), {})

        # Photo1: first animation, should be clickEffect (On Click)
        p1_ok = photo1_anim.get('node_type') in ('clickEffect', 'withEffect')
        # Photo2 and Photo3: should be afterEffect (After Previous)
        p2_ok = photo2_anim.get('node_type') == 'afterEffect'
        p3_ok = photo3_anim.get('node_type') == 'afterEffect'

        print(f"  Photo1 trigger: {photo1_anim.get('node_type')} (expect clickEffect) -> {'OK' if p1_ok else 'FAIL'}")
        print(f"  Photo2 trigger: {photo2_anim.get('node_type')} (expect afterEffect) -> {'OK' if p2_ok else 'FAIL'}")
        print(f"  Photo3 trigger: {photo3_anim.get('node_type')} (expect afterEffect) -> {'OK' if p3_ok else 'FAIL'}")

        if p1_ok and p2_ok and p3_ok:
            print(f"PASS: Component 3 — Correct animation triggers (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 3 — Incorrect triggers")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: 0.5s delay between sequential animations (0.2 points)
    try:
        parent_delays = get_parent_delays(file_path, 3)
        print(f"  Parent delays: {parent_delays}")

        # Photo2 should have 500ms delay, Photo3 should have 500ms delay
        p2_delay = parent_delays.get(str(photo2_id))
        p3_delay = parent_delays.get(str(photo3_id))

        print(f"  Photo2 delay: {p2_delay}ms (expect 500)")
        print(f"  Photo3 delay: {p3_delay}ms (expect 500)")

        if p2_delay == 500 and p3_delay == 500:
            print(f"PASS: Component 4 — Correct 0.5s delays (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 4 — Delays: Photo2={p2_delay}, Photo3={p3_delay} (expected 500)")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
