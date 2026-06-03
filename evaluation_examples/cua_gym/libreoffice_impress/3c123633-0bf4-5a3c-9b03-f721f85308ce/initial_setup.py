"""
Initial Setup: Photo Story presentation with 6 slides, 3 images on slide 4.
Task ID: impress_ma_061
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import io

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_061'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_placeholder_image(path, width, height, color, label):
    """Create a simple colored image with a label for the presentation."""
    img = Image.new('RGB', (width, height), color)
    img.save(path)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Our Weekend Photo Story"
    slide1.placeholders[1].text = "A Visual Journey Through Autumn in Vermont"

    # --- Slide 2: Section Header ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide2.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Chapter 1: Morning at the Farm"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x62)

    p2 = tf.add_paragraph()
    p2.text = "We started our day visiting the Shelburne Farms estate, where the morning mist hung low over the rolling hills."
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.runs[0]
    r2.font.size = Pt(18)
    r2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # --- Slide 3: Description slide ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox3 = slide3.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(5))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "The Covered Bridge Trail"
    p3.alignment = PP_ALIGN.LEFT
    r3 = p3.runs[0]
    r3.font.size = Pt(28)
    r3.font.bold = True
    r3.font.color.rgb = RGBColor(0x8B, 0x45, 0x13)

    for text in [
        "After the farm, we drove along Route 100 to explore the famous covered bridges.",
        "The Warren Covered Bridge, built in 1880, still carries traffic across the Mad River.",
        "Leaves in brilliant shades of red and gold framed every bridge we visited.",
        "We stopped for apple cider at a roadside stand near Waitsfield.",
    ]:
        p = tf3.add_paragraph()
        p.text = text
        p.space_before = Pt(8)
        r = p.runs[0]
        r.font.size = Pt(16)
        r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 4: Three Photos (the key slide) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])

    # Add a title text box
    txTitle = slide4.shapes.add_textbox(Inches(1), Inches(0.3), Inches(11), Inches(1))
    tf_title = txTitle.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "Highlights from the Hike"
    p_title.alignment = PP_ALIGN.CENTER
    r_title = p_title.runs[0]
    r_title.font.size = Pt(30)
    r_title.font.bold = True
    r_title.font.color.rgb = RGBColor(0x2E, 0x4A, 0x62)

    # Create 3 placeholder images
    colors = [(0x4A, 0x90, 0xD9), (0x5C, 0xB8, 0x5C), (0xD9, 0x8A, 0x4A)]
    names = ['Photo1', 'Photo2', 'Photo3']
    img_width = Inches(3.5)
    img_height = Inches(4.5)
    spacing = Inches(0.5)
    start_left = Inches(0.917)  # center three images
    top = Inches(1.8)

    for i, (color, name) in enumerate(zip(colors, names)):
        img_path = f'/home/user/_temp_{name}.png'
        create_placeholder_image(img_path, 700, 900, color, name)
        left = start_left + i * (img_width + spacing)
        pic = slide4.shapes.add_picture(img_path, int(left), int(top), int(img_width), int(img_height))
        pic.name = name
        # Clean up temp image
        os.remove(img_path)

    # --- Slide 5: More text ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox5 = slide5.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(5))
    tf5 = txBox5.text_frame
    tf5.word_wrap = True
    p5 = tf5.paragraphs[0]
    p5.text = "Evening by the Lake"
    p5.alignment = PP_ALIGN.LEFT
    r5 = p5.runs[0]
    r5.font.size = Pt(28)
    r5.font.bold = True
    r5.font.color.rgb = RGBColor(0x2E, 0x4A, 0x62)

    for text in [
        "As the sun began to set, we made our way to Lake Champlain.",
        "The water reflected the fiery orange and purple sky perfectly.",
        "A family of loons called out across the still water.",
        "We shared stories and hot cocoa as the stars appeared one by one.",
        "It was the kind of evening you never want to end.",
    ]:
        p = tf5.add_paragraph()
        p.text = text
        p.space_before = Pt(8)
        r = p.runs[0]
        r.font.size = Pt(16)
        r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 6: Closing ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox6 = slide6.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf6 = txBox6.text_frame
    tf6.word_wrap = True
    p6 = tf6.paragraphs[0]
    p6.text = "Thank You for Watching"
    p6.alignment = PP_ALIGN.CENTER
    r6 = p6.runs[0]
    r6.font.size = Pt(40)
    r6.font.bold = True
    r6.font.color.rgb = RGBColor(0x2E, 0x4A, 0x62)

    p7 = tf6.add_paragraph()
    p7.text = "Photos by Elena & David  |  October 2025"
    p7.alignment = PP_ALIGN.CENTER
    r7 = p7.runs[0]
    r7.font.size = Pt(18)
    r7.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
