"""
Reward Script: Add structured speaker notes with 'NOTES: ' prefix to all 4 slides
Task ID: osworld_impress_slide_notes_012
Domain: libreoffice_impress
Scoring:
  - Component 1: Slide 1 notes == 'NOTES: Demo opening: greet attendees and set context'  (0.25 pts)
  - Component 2: Slide 2 notes == 'NOTES: Show core feature walkthrough — 5 minutes'       (0.25 pts)
  - Component 3: Slide 3 notes == 'NOTES: Live Q&A — address top 3 questions'              (0.25 pts)
  - Component 4: Slide 4 notes == 'NOTES: Wrap up with call to action and next steps'       (0.25 pts)
  Total: 1.0
"""

import os

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_slide_notes_012'

# Expected notes per slide (task-specified, with required 'NOTES: ' prefix)
EXPECTED_NOTES = [
    "NOTES: Demo opening: greet attendees and set context",
    "NOTES: Show core feature walkthrough \u2014 5 minutes",
    "NOTES: Live Q&A \u2014 address top 3 questions",
    "NOTES: Wrap up with call to action and next steps",
]


def get_slide_notes(slide):
    """Return stripped notes text for a slide, or empty string on error."""
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Each of the 4 slides must have the exact expected note text (with NOTES: prefix).
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have exactly 4 slides
    num_slides = len(prs.slides)
    if num_slides < 4:
        print(f"CRITICAL: Expected 4 slides, found {num_slides}. Cannot score.")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide 1 notes match expected (0.25 points)
    try:
        actual = get_slide_notes(prs.slides[0])
        expected = EXPECTED_NOTES[0]
        if actual == expected:
            print(f"PASS: Component 1 — Slide 1 notes correct: {repr(actual)} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — Slide 1 notes mismatch. Expected: {repr(expected)}, Found: {repr(actual)}")
    except Exception as e:
        print(f"ERROR: Component 1 — Slide 1 notes check failed: {e}")

    # Component 2: Slide 2 notes match expected (0.25 points)
    try:
        actual = get_slide_notes(prs.slides[1])
        expected = EXPECTED_NOTES[1]
        if actual == expected:
            print(f"PASS: Component 2 — Slide 2 notes correct: {repr(actual)} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 — Slide 2 notes mismatch. Expected: {repr(expected)}, Found: {repr(actual)}")
    except Exception as e:
        print(f"ERROR: Component 2 — Slide 2 notes check failed: {e}")

    # Component 3: Slide 3 notes match expected (0.25 points)
    try:
        actual = get_slide_notes(prs.slides[2])
        expected = EXPECTED_NOTES[2]
        if actual == expected:
            print(f"PASS: Component 3 — Slide 3 notes correct: {repr(actual)} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 — Slide 3 notes mismatch. Expected: {repr(expected)}, Found: {repr(actual)}")
    except Exception as e:
        print(f"ERROR: Component 3 — Slide 3 notes check failed: {e}")

    # Component 4: Slide 4 notes match expected (0.25 points)
    try:
        actual = get_slide_notes(prs.slides[3])
        expected = EXPECTED_NOTES[3]
        if actual == expected:
            print(f"PASS: Component 4 — Slide 4 notes correct: {repr(actual)} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 4 — Slide 4 notes mismatch. Expected: {repr(expected)}, Found: {repr(actual)}")
    except Exception as e:
        print(f"ERROR: Component 4 — Slide 4 notes check failed: {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the VM environment
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
