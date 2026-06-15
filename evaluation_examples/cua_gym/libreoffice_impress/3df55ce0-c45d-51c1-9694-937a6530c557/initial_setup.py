"""
Initial Setup: Product demo presentation with 4 slides, no speaker notes
Task ID: osworld_impress_slide_notes_012
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_slide_notes_012'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # Use standard 16:9 widescreen layout
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "NovaSuite 3.0 Product Demo"
    subtitle = slide1.placeholders[1]
    subtitle.text = "Transforming Team Productivity"
    # No notes on this slide (initial state)

    # --- Slide 2: Core Features ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Core Features"
    content2 = slide2.placeholders[1]
    tf2 = content2.text_frame
    tf2.text = "Unified Dashboard"
    p2a = tf2.add_paragraph()
    p2a.text = "Real-time Analytics Engine"
    p2a.level = 1
    p2b = tf2.add_paragraph()
    p2b.text = "Automated Workflow Builder"
    p2b.level = 1
    p2c = tf2.add_paragraph()
    p2c.text = "One-click Integrations (Slack, Jira, Salesforce)"
    p2c.level = 1
    p2d = tf2.add_paragraph()
    p2d.text = "AI-Powered Smart Suggestions"
    p2d.level = 1
    # No notes on this slide (initial state)

    # --- Slide 3: Live Q&A ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Live Q&A Session"
    content3 = slide3.placeholders[1]
    tf3 = content3.text_frame
    tf3.text = "Top Questions from Attendees"
    p3a = tf3.add_paragraph()
    p3a.text = "How does NovaSuite compare to existing tools?"
    p3a.level = 1
    p3b = tf3.add_paragraph()
    p3b.text = "What is the onboarding timeline?"
    p3b.level = 1
    p3c = tf3.add_paragraph()
    p3c.text = "Is enterprise SSO supported?"
    p3c.level = 1
    # No notes on this slide (initial state)

    # --- Slide 4: Next Steps ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Next Steps"
    content4 = slide4.placeholders[1]
    tf4 = content4.text_frame
    tf4.text = "Get Started Today"
    p4a = tf4.add_paragraph()
    p4a.text = "Schedule a personalized onboarding call"
    p4a.level = 1
    p4b = tf4.add_paragraph()
    p4b.text = "Start your 30-day free trial at novasuite.io/trial"
    p4b.level = 1
    p4c = tf4.add_paragraph()
    p4c.text = "Contact our sales team: sales@novasuite.io"
    p4c.level = 1
    # No notes on this slide (initial state)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
