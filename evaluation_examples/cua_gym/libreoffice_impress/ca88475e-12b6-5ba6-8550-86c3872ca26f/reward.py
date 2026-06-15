"""
FINAL REWARD SCRIPT - SUCCESS
Task: I'm working on a presentation and need to make navigation easier for the audience. How can I link 'Details' on slide 2 directly to slide 5 in LibreOffice Impress?
Generated: 2025-08-07 13:30:18
Status: success
Model: o4-mini
Total Steps: 9
"""

import os
from pptx import Presentation
import zipfile
import xml.etree.ElementTree as ET

def verify_impress_hyperlink(file_path, source_slide_idx, link_text, target_slide_idx):
    print("Checking hyperlink task in Impress...")
    total_score = 0.0
    max_score = 1.0

    # 1. File exists check (0.15)
    print(f"1. Checking if file exists: {file_path}")
    if os.path.exists(file_path):
        print("✓ File exists (0.15)")
        total_score += 0.15
    else:
        print("✗ File not found (0 points)")
        print(f"REWARD: {total_score}")
        return

    # 2. Load presentation (0.15)
    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded with {len(prs.slides)} slides (0.15)")
        total_score += 0.15
    except Exception as e:
        print(f"✗ Failed to load presentation: {e} (0 points)")
        print(f"REWARD: {total_score}")
        return

    # 3. Slide count check (0.15)
    needed = max(source_slide_idx+1, target_slide_idx+1)
    actual = len(prs.slides)
    print(f"2. Checking slide count: need ≥{needed}, found {actual}")
    if actual >= needed:
        print("✓ Sufficient number of slides (0.15)")
        total_score += 0.15
    else:
        print("✗ Not enough slides (0 points)")
        print(f"REWARD: {total_score}")
        return

    # 4. Find shape with link text (0.2)
    print(f"3. Searching for shape containing text '{link_text}' on slide {source_slide_idx+1}")
    slide = prs.slides[source_slide_idx]
    shape_found = False
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text and link_text.lower() in shape.text.lower():
            print(f"✓ Found shape with text including '{link_text}' (0.2)")
            total_score += 0.2
            shape_found = True
            break
    if not shape_found:
        print(f"✗ No shape containing text '{link_text}' found (0 points)")

    # 5. Check slide relationships for target slide hyperlink (0.35)
    rels_path = f"ppt/slides/_rels/slide{source_slide_idx+1}.xml.rels"
    print(f"4. Inspecting relationships file '{rels_path}' for link to slide {target_slide_idx+1}")
    try:
        with zipfile.ZipFile(file_path) as z:
            if rels_path in z.namelist():
                rels_xml = z.read(rels_path)
                ns = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}
                root = ET.fromstring(rels_xml)
                found_link = False
                for rel in root.findall('r:Relationship', ns):
                    typ = rel.attrib.get('Type', '')
                    target = rel.attrib.get('Target', '')
                    if typ.endswith('/slide') and target == f"slide{target_slide_idx+1}.xml":
                        print(f"✓ Relationship to slide {target_slide_idx+1} found (0.35)")
                        total_score += 0.35
                        found_link = True
                        break
                if not found_link:
                    print(f"✗ No relationship linking to slide {target_slide_idx+1} (0 points)")
            else:
                print(f"✗ Relationships file not found for slide {source_slide_idx+1} (0 points)")
    except Exception as e:
        print(f"✗ Error reading relationships: {e} (0 points)")

    # Final score rounding to avoid floating point errors
    final_score = min(total_score, max_score)
    if abs(final_score - max_score) < 1e-6:
        final_score = max_score
    else:
        final_score = round(final_score, 3)

    print(f"Final score: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score

if __name__ == '__main__':
    file_path = '/home/user/im_working_on_a_presentation_and_need_to_make_navigation_easier_for_the_audience_how_can_i_link_deta.pptx'
    verify_impress_hyperlink(file_path, source_slide_idx=1, link_text='Details', target_slide_idx=4)
