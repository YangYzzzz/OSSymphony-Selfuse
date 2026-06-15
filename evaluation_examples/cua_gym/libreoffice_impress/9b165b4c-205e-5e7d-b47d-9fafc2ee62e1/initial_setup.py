"""
Initial Setup: Digital Kiosk presentation with 10 slides, no automatic timing
Task ID: impress_gf2_008
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_008'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()

    # --- Slide 1: Welcome / Title ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "NovaTech Digital Experience"
    slide1.placeholders[1].text = "Innovation Hub Interactive Kiosk"
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)

    # --- Slide 2: Company Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "About NovaTech"
    slide2.placeholders[1].text = (
        "Founded in 2018, NovaTech Solutions delivers cutting-edge digital "
        "infrastructure for smart cities and enterprise environments.\n\n"
        "Headquarters: Austin, TX\n"
        "Employees: 1,200+\n"
        "Global offices: 8 countries"
    )

    # --- Slide 3: Mission Statement ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Our Mission"
    slide3.placeholders[1].text = (
        "Empowering organizations to create seamless, intelligent environments "
        "through integrated IoT platforms and human-centered design.\n\n"
        "Core Values:\n"
        "- Innovation through collaboration\n"
        "- Sustainable technology solutions\n"
        "- Customer-first engineering"
    )

    # --- Slide 4: Smart Building Solutions ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Smart Building Solutions"
    slide4.placeholders[1].text = (
        "Automated Climate Control\n"
        "Real-time occupancy sensing reduces energy costs by 35%.\n\n"
        "Predictive Maintenance\n"
        "ML-driven alerts prevent 92% of critical equipment failures.\n\n"
        "Integrated Access Management\n"
        "Biometric and RFID systems supporting 50,000+ daily authentications."
    )

    # --- Slide 5: Digital Signage Platform ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Digital Signage Platform"
    slide5.placeholders[1].text = (
        "CloudCanvas v3.2 — Enterprise Content Management\n\n"
        "Features:\n"
        "- Remote content scheduling across 10,000+ displays\n"
        "- Real-time analytics dashboard with audience metrics\n"
        "- Multi-zone layout engine with 4K HDR support\n"
        "- Emergency broadcast override capability"
    )

    # --- Slide 6: IoT Sensor Network ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "IoT Sensor Network"
    slide6.placeholders[1].text = (
        "SensorMesh Architecture\n\n"
        "Coverage: 2.4M active sensors deployed globally\n"
        "Latency: Sub-100ms edge processing\n"
        "Protocols: MQTT, CoAP, LoRaWAN, Zigbee 3.0\n"
        "Uptime: 99.97% availability SLA\n\n"
        "Applications: environmental monitoring, asset tracking, "
        "space utilization, air quality management"
    )

    # --- Slide 7: Case Study — Metro City Transit ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Case Study: Metro City Transit"
    slide7.placeholders[1].text = (
        "Challenge:\n"
        "Modernize 340 transit stations with real-time passenger information.\n\n"
        "Solution:\n"
        "Deployed 2,800 interactive kiosks with route planning, "
        "live departure boards, and accessibility features.\n\n"
        "Results:\n"
        "- 28% increase in passenger satisfaction scores\n"
        "- 15% reduction in customer service inquiries\n"
        "- $4.2M annual savings in printed materials"
    )

    # --- Slide 8: Product Roadmap 2026 ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Product Roadmap 2026"
    slide8.placeholders[1].text = (
        "Q1: AI-powered content personalization engine\n"
        "Q2: Augmented reality wayfinding integration\n"
        "Q3: Gesture and voice interaction modules\n"
        "Q4: Federated analytics platform launch\n\n"
        "Investment: $18M R&D budget allocated\n"
        "Partnerships: Google Cloud, NVIDIA, Qualcomm"
    )

    # --- Slide 9: Get In Touch ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Get In Touch"
    slide9.placeholders[1].text = (
        "Sales Inquiries: sales@novatech.io\n"
        "Technical Support: support@novatech.io\n"
        "Phone: +1 (512) 555-0147\n\n"
        "Visit our demo center:\n"
        "NovaTech Innovation Hub\n"
        "2400 Technology Blvd, Austin, TX 78759"
    )

    # --- Slide 10: Thank You ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[0])
    slide10.shapes.title.text = "Thank You"
    slide10.placeholders[1].text = "www.novatech.io"
    fill = slide10.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
