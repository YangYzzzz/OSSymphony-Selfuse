"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 187’s heading is kind of blending into the background—how do I put a #FFFF00 character highlight behind that title text in LibreOffice Impress?
Generated: 2025-09-10 18:13:11
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
import zipfile
import re
from pptx import Presentation


def verify_slide_highlight(file_path: str,
                           slide_number: int = 187,
                           expected_color: str = "FFFF00") -> float:
    """Verify that a specific slide has text highlighted with the
    expected RGB hex color (default #FFFF00).

    Scoring rubric (progressive up to 1.0):
      • 0.1  – slide exists
      • 0.1  – slide contains text
      • 0.3  – <a:highlight> element present in slide XML
      • 0.5  – highlight colour matches expected value
    Returns a float between 0.0 and 1.0 inclusive.
    """

    print(f"Verifying highlight on slide {slide_number} in: {file_path}")
    max_score = 1.0
    score = 0.0

    # ---------- 1. File & slide checks (small credit) ----------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0

    try:
        prs = Presentation(file_path)
        total_slides = len(prs.slides)
        print(f"Loaded presentation with {total_slides} slides")
    except Exception as e:
        print(f"✗ Unable to load presentation: {e}")
        return 0.0

    if slide_number > total_slides:
        print(f"✗ Slide {slide_number} does not exist (only {total_slides} slides)")
        return 0.0
    else:
        print("✓ Slide exists")
        score += 0.1  # small, because main requirement is highlight

    # Check that the slide actually contains some text
    slide = prs.slides[slide_number - 1]
    has_text = any(getattr(shape, "text", "").strip() for shape in slide.shapes)
    if has_text:
        print("✓ Text content found on slide")
        score += 0.1
    else:
        print("✗ No text content found on slide")

    # ---------- 2. Deep XML inspection for highlight ----------
    slide_xml_path = f"ppt/slides/slide{slide_number}.xml"
    try:
        with zipfile.ZipFile(file_path, "r") as pptx_zip:
            if slide_xml_path not in pptx_zip.namelist():
                print(f"✗ Slide XML {slide_xml_path} not found in archive")
                return score
            xml_bytes = pptx_zip.read(slide_xml_path)
    except Exception as e:
        print(f"✗ Error reading slide XML: {e}")
        return score

    xml_text = xml_bytes.decode("utf-8", errors="ignore")

    # Look for <a:highlight> element
    if "<a:highlight" in xml_text:
        print("✓ <a:highlight> element detected")
        score += 0.3
    else:
        print("✗ No <a:highlight> element detected on slide")
        return min(score, max_score)

    # Verify the highlight colour value
    colour_regex = re.compile(r"val=\"?([0-9A-Fa-f]{6})\"?")
    colour_matches = colour_regex.findall(xml_text)

    colour_found = any(match.upper() == expected_color.upper() for match in colour_matches)
    if colour_found:
        print(f"✓ Highlight colour {expected_color} found")
        score += 0.5
    else:
        print(f"✗ Highlight colour {expected_color} not found")

    final_score = min(score, max_score)
    print(f"Total score: {final_score}")
    return final_score


if __name__ == "__main__":
    target_file = "/home/user/slide_187s_heading_is_kind_of_blending_into_the_backgroundhow_do_i_put_a_ffff00_character_highlight__golden.pptx"
    reward = verify_slide_highlight(target_file)
    print(f"REWARD: {reward}")
