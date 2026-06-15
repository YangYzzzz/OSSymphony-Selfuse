"""
Initial Setup: Insert a table on slide 4 of a presentation
Task ID: impress_gf5_006
Domain: libreoffice_impress

Creates an 8-slide business presentation. Slide 4 ("Revenue Comparison")
has a title but NO table — the agent must add the table.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_006'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_slide(prs, layout_idx, title_text, body_lines):
    """Helper to add a slide with title and bullet-point body text."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title_text
    # Find body placeholder (index 1)
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide (layout 0) ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Global Sales Report 2025"
    slide1.placeholders[1].text = "Prepared by the Finance & Strategy Team\nMarch 2025"

    # --- Slide 2: Executive Summary (layout 1 = Title + Content) ---
    add_text_slide(prs, 1, "Executive Summary", [
        "Total revenue grew 18% year-over-year to $142.3M",
        "North America remained the largest market at 45% share",
        "Three new product lines contributed $23.7M in first-year sales",
        "Operating margins improved from 22.1% to 25.8%",
        "Customer retention rate held steady at 91.4%",
    ])

    # --- Slide 3: Market Analysis (layout 1) ---
    add_text_slide(prs, 1, "Market Analysis", [
        "Enterprise segment: $67.2M (+21% YoY)",
        "Mid-market segment: $48.5M (+15% YoY)",
        "SMB segment: $26.6M (+12% YoY)",
        "Competitive pressure increasing in APAC region",
        "Digital transformation driving demand across all verticals",
    ])

    # --- Slide 4: Revenue Comparison — EMPTY content area, NO TABLE ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Revenue Comparison"
    # Clear the body placeholder to leave empty content area
    body = slide4.placeholders[1]
    tf = body.text_frame
    tf.clear()
    tf.paragraphs[0].text = ""

    # --- Slide 5: Regional Performance (layout 1) ---
    add_text_slide(prs, 1, "Regional Performance", [
        "North America: $64.0M — strongest in enterprise SaaS",
        "Europe: $38.9M — growth driven by GDPR compliance tools",
        "Asia-Pacific: $27.1M — expanding partner network in Japan and Australia",
        "Latin America: $8.5M — new office in São Paulo operational since Q2",
        "Middle East & Africa: $3.8M — early-stage market entry",
    ])

    # --- Slide 6: Customer Insights (layout 1) ---
    add_text_slide(prs, 1, "Customer Insights", [
        "Net Promoter Score improved from 62 to 71",
        "Average deal size increased 14% to $87,400",
        "Time-to-close reduced by 8 days on average",
        "Top requested feature: advanced analytics dashboard",
        "Customer success team expanded to 45 specialists",
    ])

    # --- Slide 7: Growth Strategy (layout 1) ---
    add_text_slide(prs, 1, "Growth Strategy", [
        "Launch AI-powered forecasting module in Q3 2025",
        "Expand APAC sales team by 30% by end of year",
        "Target 3 strategic acquisitions in adjacent markets",
        "Invest $12M in R&D for next-generation platform",
        "Establish dedicated vertical solutions for healthcare and fintech",
    ])

    # --- Slide 8: Next Steps (layout 1) ---
    add_text_slide(prs, 1, "Next Steps", [
        "Finalize Q2 product roadmap by April 15",
        "Complete Series C funding round — target $50M",
        "Onboard 12 new enterprise accounts in pipeline",
        "Roll out updated pricing model for mid-market tier",
        "Schedule board review for June 2025",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
