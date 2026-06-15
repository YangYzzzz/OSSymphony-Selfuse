"""
Reward Script: Insert a 4x3 table on slide 4 with headers and data
Task ID: impress_gf5_006
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Table exists on slide 4 with 4 rows and 3 columns
  Component 2 (0.30): Header row contains 'Product', 'Q1 Revenue', 'Q2 Revenue'
  Component 3 (0.25): Data rows 2-4 contain product names and revenue figures
  Component 4 (0.15): Header row is visually distinct (bold and/or colored text)
"""

import os
import sys

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_006'


def persist_app_state(domain: str):
    """Attempt to save any unsaved changes in LibreOffice."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file has at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Precondition — need at least 4 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)

    # Find table shape on slide 4
    table_shape = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_shape = shape
            break

    # Component 1: Table exists on slide 4 with 4 rows and 3 columns (0.30 points)
    try:
        if table_shape is None:
            print("FAIL: Component 1 — No table found on slide 4")
        else:
            table = table_shape.table
            num_rows = len(table.rows)
            num_cols = len(table.columns)
            if num_rows == 4 and num_cols == 3:
                print(f"PASS: Component 1 — Table found on slide 4 with {num_rows}x{num_cols} dimensions (0.30 pts)")
                total_score += 0.30
            else:
                print(f"FAIL: Component 1 — Table dimensions are {num_rows}x{num_cols}, expected 4x3")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # If no table, remaining components cannot be checked
    if table_shape is None:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    table = table_shape.table

    # Component 2: Header row has correct text (0.30 points)
    try:
        expected_headers = ['Product', 'Q1 Revenue', 'Q2 Revenue']
        actual_headers = []
        for c in range(min(len(table.columns), 3)):
            actual_headers.append(table.cell(0, c).text.strip())

        matches = 0
        for i, (expected, actual) in enumerate(zip(expected_headers, actual_headers)):
            if expected.lower() == actual.lower():
                matches += 1
            else:
                print(f"  Header mismatch col {i}: expected '{expected}', got '{actual}'")

        if matches == 3:
            print(f"PASS: Component 2 — All 3 headers match: {actual_headers} (0.30 pts)")
            total_score += 0.30
        elif matches >= 2:
            partial = 0.20
            print(f"PARTIAL: Component 2 — {matches}/3 headers match: {actual_headers} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — Headers found: {actual_headers}, expected: {expected_headers}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Data rows contain product names and revenue figures (0.25 points)
    try:
        filled_data_rows = 0
        for r in range(1, min(len(table.rows), 4)):
            product_name = table.cell(r, 0).text.strip()
            q1_rev = table.cell(r, 1).text.strip()
            q2_rev = table.cell(r, 2).text.strip()

            # Product name must be non-empty string
            has_product = len(product_name) > 0
            # Revenue values should contain digits (e.g., "$1,245,000" or "1245000")
            has_q1 = any(ch.isdigit() for ch in q1_rev) if q1_rev else False
            has_q2 = any(ch.isdigit() for ch in q2_rev) if q2_rev else False

            if has_product and has_q1 and has_q2:
                filled_data_rows += 1
                print(f"  Row {r}: '{product_name}' | '{q1_rev}' | '{q2_rev}' — OK")
            else:
                print(f"  Row {r}: '{product_name}' | '{q1_rev}' | '{q2_rev}' — incomplete")

        if filled_data_rows == 3:
            print(f"PASS: Component 3 — All 3 data rows have product names and revenue figures (0.25 pts)")
            total_score += 0.25
        elif filled_data_rows >= 2:
            partial = 0.15
            print(f"PARTIAL: Component 3 — {filled_data_rows}/3 data rows complete ({partial} pts)")
            total_score += partial
        elif filled_data_rows >= 1:
            partial = 0.08
            print(f"PARTIAL: Component 3 — {filled_data_rows}/3 data rows complete ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No data rows have complete product/revenue data")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Header row is visually distinct — bold and/or different color (0.15 points)
    try:
        header_distinct = False
        for c in range(min(len(table.columns), 3)):
            cell = table.cell(0, c)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    # Check bold
                    if run.font.bold is True:
                        header_distinct = True
                        break
                    # Check font color is different from default black
                    try:
                        if run.font.color.type is not None and run.font.color.rgb is not None:
                            color_str = str(run.font.color.rgb)
                            if color_str != '000000':
                                header_distinct = True
                                break
                    except:
                        pass
                if header_distinct:
                    break
            if header_distinct:
                break

        # Also check for header row background color via XML
        if not header_distinct:
            try:
                tc = table.cell(0, 0)._tc
                tcPr = tc.find(qn('a:tcPr'))
                if tcPr is not None:
                    solid = tcPr.find(qn('a:solidFill'))
                    if solid is not None:
                        header_distinct = True
            except:
                pass

        # Also check table style flags (firstRow banding)
        if not header_distinct:
            try:
                tblPr = table._tbl.find(qn('a:tblPr'))
                if tblPr is not None:
                    first_row = tblPr.get('firstRow', '0')
                    if first_row == '1':
                        header_distinct = True
            except:
                pass

        if header_distinct:
            print(f"PASS: Component 4 — Header row is visually distinct (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 — Header row is not visually distinct from data rows")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
