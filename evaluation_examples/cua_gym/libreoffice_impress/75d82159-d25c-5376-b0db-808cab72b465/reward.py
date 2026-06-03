"""
FINAL REWARD SCRIPT - SUCCESS
Task: Convert straight quotes to typographic quotes across the document.
Generated: 2025-10-17 10:42:01
Status: success
Model: azure-o3
Total Steps: 1
"""

from pptx import Presentation
import os

def collect_texts(presentation):
    """Collect all text strings from shapes and notes in the presentation."""
    texts = []

    def extract_from_shape(shape):
        # Recursively extract text from grouped shapes as well
        if shape.shape_type == 6:  # group shape
            for shp in shape.shapes:
                extract_from_shape(shp)
        else:
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                texts.append(shape.text)

    for slide in presentation.slides:
        for shape in slide.shapes:
            extract_from_shape(shape)
        # Also check speaker notes (if any)
        try:
            notes_slide = slide.notes_slide
            if notes_slide and notes_slide.notes_text_frame:
                texts.append(notes_slide.notes_text_frame.text)
        except Exception:
            pass  # Some slides may not have notes (ignore)

    return texts

def verify_typographic_quotes(file_path):
    """Verify that straight quotes have been replaced by typographic (curly) quotes."""
    print(f"Verifying typographic quotes for file: {file_path}")

    # Preliminary checks (no points awarded for these)
    if not os.path.exists(file_path):
        print("✗ File not found")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        return 0.0

    # Gather all text elements
    texts = collect_texts(prs)
    total_text_count = len(texts)
    print(f"Collected {total_text_count} text elements from presentation")

    # Flags for various quote types
    ascii_single_found = False  # Straight single quote '
    ascii_double_found = False  # Straight double quote "
    curly_found = False         # Any curly quote character

    curly_chars = {"\u2018", "\u2019", "\u201C", "\u201D"}

    for text in texts:
        if text is None:
            continue
        # Straight quotes detection
        if "'" in text:
            ascii_single_found = True
        if '"' in text:
            ascii_double_found = True
        # Curly quotes detection
        if any(c in text for c in curly_chars):
            curly_found = True

    # Progressive scoring
    score = 0.0

    # 0.4 points if no straight double quotes remain
    if not ascii_double_found:
        print("✓ No straight double quotes found (0.4)")
        score += 0.4
    else:
        print("✗ Straight double quotes still present")

    # 0.4 points if no straight single quotes remain
    if not ascii_single_found:
        print("✓ No straight single quotes found (0.4)")
        score += 0.4
    else:
        print("✗ Straight single quotes still present")

    # 0.2 points if at least one curly quote character is present somewhere in the document
    if curly_found:
        print("✓ Curly quotation marks present (0.2)")
        score += 0.2
    else:
        print("✗ No curly quotation marks detected")

    # Ensure score does not exceed 1.0
    final_score = min(score, 1.0)
    print(f"REWARD: {final_score}")

    return final_score

# -----------------------------------------------------------------------------
# ACTUAL EXECUTION
# -----------------------------------------------------------------------------
if __name__ == "__main__":
    file_path = "/home/user/convert_straight_quotes_to_typographic_quotes_across_the_document.pptx"
    verify_typographic_quotes(file_path)

