"""
Initial Setup: Conference talk presentation with 5 slides, all manual advance.
Task ID: impress_gf2_017
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import zipfile
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_017'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Scaling Microservices in Production"
    slide1.placeholders[1].text = "Dr. Elena Vasquez — Platform Engineering Summit 2026"

    # --- Slide 2: Agenda / Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Current Architecture Challenges"
    items2 = [
        "Container Orchestration Patterns",
        "Service Mesh Implementation",
        "Observability at Scale",
        "Performance Benchmarks & Results",
    ]
    for item in items2:
        p = body2.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 3: Architecture Diagram (text description) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Current Architecture Overview"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "API Gateway → Load Balancer → Service Cluster"
    details3 = [
        "12 core services handling 45K RPS peak traffic",
        "PostgreSQL + Redis caching layer (95% hit rate)",
        "Event-driven messaging via Apache Kafka",
        "Kubernetes cluster: 48 nodes across 3 availability zones",
        "Average response latency: 23ms (p50), 87ms (p99)",
    ]
    for d in details3:
        p = body3.add_paragraph()
        p.text = d
        p.level = 1

    # --- Slide 4: Performance Results ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Performance Benchmarks"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Key Metrics After Migration"
    metrics = [
        "Throughput increased from 12K to 45K RPS (+275%)",
        "P99 latency reduced from 340ms to 87ms (-74%)",
        "Infrastructure cost reduced by $18,400/month (-32%)",
        "Deployment frequency: 3x/week → 12x/day",
        "Mean time to recovery: 45 min → 4.2 min",
    ]
    for m in metrics:
        p = body4.add_paragraph()
        p.text = m
        p.level = 0

    # --- Slide 5: Q&A / Closing ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[0])
    slide5.shapes.title.text = "Questions & Discussion"
    slide5.placeholders[1].text = "elena.vasquez@techcorp.io  |  @elenavdev  |  github.com/evasquez"

    # Save the presentation
    prs.save(OUTPUT)

    # Now strip any default transition/advance attributes via XML to ensure manual advance only
    _ensure_manual_advance(OUTPUT)

    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


def _ensure_manual_advance(pptx_path):
    """Remove any automatic advance timing from all slides via XML manipulation."""
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    tmp_path = pptx_path + '.tmp'

    with zipfile.ZipFile(pptx_path, 'r') as zin, zipfile.ZipFile(tmp_path, 'w') as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith('ppt/slides/slide') and item.filename.endswith('.xml'):
                root = ET.fromstring(data)
                # Find and remove advTm (auto-advance) from any transition element
                for tr in root.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}transition'):
                    if 'advTm' in tr.attrib:
                        del tr.attrib['advTm']
                    # Ensure advClick is not set to "0" (which would disable click advance)
                    tr.attrib['advClick'] = '1'
                data = ET.tostring(root, xml_declaration=True, encoding='UTF-8')
            zout.writestr(item, data)

    os.replace(tmp_path, pptx_path)


create_initial()
