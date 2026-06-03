"""
Initial Setup: Create Tech_Architecture.pptx with 8 slides, slide 6 has title 'Technology Stack' but no layer diagram.
Task ID: impress_ps_032
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_032'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.text = body_lines[0]
    for line in body_lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
    return slide


def add_title_only_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    add_title_slide(prs, "Tech Architecture Overview",
                    "Cloud-Native Platform Design\nQ1 2025 Technical Review")

    # Slide 2: Agenda
    add_content_slide(prs, "Agenda", [
        "System Architecture Goals",
        "Microservices Design Patterns",
        "Infrastructure & Deployment",
        "Technology Stack",
        "Security & Compliance",
        "Performance Benchmarks",
        "Roadmap & Next Steps",
    ])

    # Slide 3: System Architecture Goals
    add_content_slide(prs, "System Architecture Goals", [
        "Achieve 99.99% uptime SLA for all production services",
        "Support horizontal scaling to 10M concurrent users",
        "Reduce mean deployment time to under 15 minutes",
        "Implement zero-trust security model across all services",
        "Maintain P99 latency below 200ms for API endpoints",
    ])

    # Slide 4: Microservices Design Patterns
    add_content_slide(prs, "Microservices Design Patterns", [
        "Event-driven architecture using Apache Kafka",
        "Service mesh with Istio for inter-service communication",
        "Circuit breaker pattern via Resilience4j",
        "CQRS for read/write separation on high-traffic services",
        "Saga pattern for distributed transaction management",
    ])

    # Slide 5: Infrastructure & Deployment
    add_content_slide(prs, "Infrastructure & Deployment", [
        "Multi-region AWS deployment across us-east-1 and eu-west-1",
        "Kubernetes clusters managed via EKS with Karpenter autoscaling",
        "GitOps workflow with ArgoCD for continuous deployment",
        "Terraform modules for infrastructure-as-code provisioning",
        "Prometheus + Grafana for observability and alerting",
    ])

    # Slide 6: Technology Stack - title only, empty content area (task target)
    slide6 = add_title_only_slide(prs, "Technology Stack")
    # Intentionally left empty - no shapes, no layer diagram

    # Slide 7: Security & Compliance
    add_content_slide(prs, "Security & Compliance", [
        "OAuth 2.0 / OIDC authentication via Keycloak",
        "mTLS for all service-to-service communication",
        "Automated vulnerability scanning with Trivy and Snyk",
        "SOC 2 Type II certification in progress",
        "Data encryption at rest (AES-256) and in transit (TLS 1.3)",
    ])

    # Slide 8: Roadmap & Next Steps
    add_content_slide(prs, "Roadmap & Next Steps", [
        "Q2: Migrate remaining monolith services to microservices",
        "Q2: Implement edge caching with CloudFront + Lambda@Edge",
        "Q3: Deploy ML inference pipeline on SageMaker",
        "Q3: Achieve SOC 2 Type II certification",
        "Q4: Expand to Asia-Pacific region (ap-southeast-1)",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
