"""
Reward Script: Layered architecture slide with 4 stacked rectangles on slide 6
Task ID: impress_ps_032
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): 4 rectangle shapes exist on slide 6
  Component 2 (0.25): Correct layer text labels on rectangles
  Component 3 (0.25): Correct fill colors (dark to light blue gradient)
  Component 4 (0.25): Correct stacking order (bottom-to-top) with overlap
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_032'

# Expected layer definitions: (text_substring, fill_color_hex)
# Ordered bottom-to-top (highest y-position first = bottom of slide)
EXPECTED_LAYERS = [
    ('Infrastructure (AWS)', '1A237E'),
    ('Platform (Kubernetes)', '1565C0'),
    ('Application (Node.js)', '42A5F5'),
    ('Client (React)', 'BBDEFB'),
]


def persist_app_state(domain):
    """Save any unsaved GUI state before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: python-pptx not available: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[5]  # slide 6 (0-indexed)

    # Collect all AUTO_SHAPE (rectangle) shapes on slide 6
    # Exclude placeholders and text boxes which are pre-existing
    rect_shapes = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            rect_shapes.append(shape)

    # Component 1: Exactly 4 rectangle shapes exist on slide 6 (0.25 points)
    try:
        num_rects = len(rect_shapes)
        if num_rects == 4:
            print(f"PASS: Component 1 -- Found exactly 4 rectangle shapes on slide 6 (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 -- Expected 4 rectangles on slide 6, found {num_rects}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    if len(rect_shapes) < 4:
        # Cannot proceed with detailed checks if fewer than 4 rectangles
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Sort rectangles by top position descending (bottom of slide = highest y = first)
    rect_shapes_sorted = sorted(rect_shapes, key=lambda s: s.top, reverse=True)

    # Component 2: Correct layer text labels (0.25 points)
    # Each rectangle should contain the correct layer name text
    try:
        text_matches = 0
        for i, (expected_text, _) in enumerate(EXPECTED_LAYERS):
            if i < len(rect_shapes_sorted):
                shape = rect_shapes_sorted[i]
                actual_text = shape.text.strip() if hasattr(shape, 'text') else ''
                if expected_text.lower() in actual_text.lower():
                    text_matches += 1
                    print(f"  PASS: Layer {i+1} text matches: '{actual_text}' contains '{expected_text}'")
                else:
                    print(f"  FAIL: Layer {i+1} expected text containing '{expected_text}', found '{actual_text}'")

        if text_matches == 4:
            print(f"PASS: Component 2 -- All 4 layer labels correct (0.25 pts)")
            total_score += 0.25
        elif text_matches >= 2:
            partial = 0.25 * (text_matches / 4)
            print(f"PARTIAL: Component 2 -- {text_matches}/4 labels correct ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 -- Only {text_matches}/4 labels correct")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Correct fill colors (0.25 points)
    # Each rectangle should have the correct blue shade fill
    try:
        color_matches = 0
        for i, (_, expected_color) in enumerate(EXPECTED_LAYERS):
            if i < len(rect_shapes_sorted):
                shape = rect_shapes_sorted[i]
                try:
                    fill = shape.fill
                    if fill.type == 1:  # SOLID fill
                        actual_color = str(fill.fore_color.rgb).upper()
                        expected_upper = expected_color.upper()
                        if actual_color == expected_upper:
                            color_matches += 1
                            print(f"  PASS: Layer {i+1} fill color #{actual_color} matches expected #{expected_upper}")
                        else:
                            print(f"  FAIL: Layer {i+1} fill color #{actual_color}, expected #{expected_upper}")
                    else:
                        print(f"  FAIL: Layer {i+1} fill type is {fill.type}, expected SOLID (1)")
                except Exception as e:
                    print(f"  ERROR: Layer {i+1} fill check: {e}")

        if color_matches == 4:
            print(f"PASS: Component 3 -- All 4 fill colors correct (0.25 pts)")
            total_score += 0.25
        elif color_matches >= 2:
            partial = 0.25 * (color_matches / 4)
            print(f"PARTIAL: Component 3 -- {color_matches}/4 colors correct ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 -- Only {color_matches}/4 colors correct")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Correct stacking order with overlap (0.25 points)
    # Rectangles should be stacked vertically with bottom layer at highest y-position
    # and each layer should slightly overlap the one below (gap < height)
    try:
        order_ok_count = 0
        overlap_ok_count = 0
        pairs_to_check = len(rect_shapes_sorted) - 1  # should be 3

        # Check that they are in descending y-order (already sorted that way)
        # Verify: each successive shape (going up) has a lower y (top) value
        for i in range(pairs_to_check):
            lower_shape = rect_shapes_sorted[i]
            upper_shape = rect_shapes_sorted[i + 1]

            # The upper shape should have a smaller top value
            if upper_shape.top < lower_shape.top:
                order_ok_count += 1
            else:
                print(f"  FAIL: Stacking order wrong: shape at y={upper_shape.top} should be above shape at y={lower_shape.top}")

            # Check overlap: the top of the lower shape should be less than
            # (top + height) of the upper shape, i.e., they overlap
            upper_bottom = upper_shape.top + upper_shape.height
            if upper_bottom > lower_shape.top:
                overlap_ok_count += 1
                print(f"  PASS: Overlap detected between layers {i+1} and {i+2}: upper_bottom={upper_bottom} > lower_top={lower_shape.top}")
            else:
                print(f"  FAIL: No overlap between layers {i+1} and {i+2}: upper_bottom={upper_bottom} <= lower_top={lower_shape.top}")

        # Also check all rectangles have equal width (structural requirement)
        widths = [s.width for s in rect_shapes_sorted]
        equal_width = len(set(widths)) == 1

        if order_ok_count == pairs_to_check and overlap_ok_count == pairs_to_check and equal_width:
            print(f"PASS: Component 4 -- Correct stacking order with overlap, equal widths (0.25 pts)")
            total_score += 0.25
        elif order_ok_count == pairs_to_check and equal_width:
            print(f"PARTIAL: Component 4 -- Order correct but overlap not verified (0.15 pts)")
            total_score += 0.15
        elif order_ok_count == pairs_to_check:
            print(f"PARTIAL: Component 4 -- Order correct but widths vary (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 4 -- Stacking order incorrect")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
