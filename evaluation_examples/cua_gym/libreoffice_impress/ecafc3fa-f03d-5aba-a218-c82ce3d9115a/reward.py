"""
FINAL REWARD SCRIPT - SUCCESS
Task: Replace straight quotes with typographic quotes for the whole document.
Generated: 2025-10-17 08:29:24
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation

# Constants for quote characters
STRAIGHT_SINGLE = "'"
STRAIGHT_DOUBLE = '"'
CURVY_SINGLES = ['‘', '’']
CURVY_DOUBLES = ['“', '”']

# Shape type constants from python-pptx (hard-coded to avoid extra import)
MSO_SHAPE_TYPE_GROUP = 6   # GROUP
MSO_SHAPE_TYPE_TABLE = 19  # TABLE


def extract_text_from_shape(shape):
    """Recursively collect text from any shape, including groups and tables."""
    texts = []

    # Basic text frames (text boxes, placeholders, etc.)
    if getattr(shape, 'has_text_frame', False):
        texts.append(shape.text_frame.text)

    # Tables: iterate through every cell
    if shape.shape_type == MSO_SHAPE_TYPE_TABLE:
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                if cell.text_frame is not None:
                    texts.append(cell.text_frame.text)

    # Group shapes: recurse into child shapes
    if shape.shape_type == MSO_SHAPE_TYPE_GROUP:
        for shp in shape.shapes:
            texts.extend(extract_text_from_shape(shp))

    return texts


def collect_all_text(presentation):
    """Return a list with every piece of text found in the presentation."""
    all_texts = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            all_texts.extend(extract_text_from_shape(shape))
    return all_texts


def count_quotes(texts):
    """Count straight and curly quote characters in a list of strings."""
    straight_single = straight_double = 0
    curly_single = curly_double = 0

    for text in texts:
        straight_single += text.count(STRAIGHT_SINGLE)
        straight_double += text.count(STRAIGHT_DOUBLE)
        curly_single += sum(text.count(c) for c in CURVY_SINGLES)
        curly_double += sum(text.count(c) for c in CURVY_DOUBLES)

    return straight_single, straight_double, curly_single, curly_double


def verify_typographic_quotes(file_path):
    """
    Reward-based verification:
    1. 0.4 pts – No straight double quotes remain
    2. 0.4 pts – No straight single quotes remain
    3. 0.2 pts – At least one curly quote exists (ensures replacement, not deletion)
    Returns a float between 0.0 and 1.0.
    """
    print(f"Verifying typographic quotes in: {file_path}")
    score = 0.0
    max_score = 1.0

    # --- Prerequisite: file must exist and load (no points for this) ---
    if not os.path.exists(file_path):
        print("✗ File not found.")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded presentation with {len(prs.slides)} slides (prerequisite – 0 points)")
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # --- Core verification ---
    texts = collect_all_text(prs)
    ss, sd, cs, cd = count_quotes(texts)

    print(f"Straight single quotes ('): {ss}")
    print(f"Straight double quotes (\"): {sd}")
    print(f"Curly single quotes: {cs}")
    print(f"Curly double quotes: {cd}")

    # A) Straight double quotes gone
    if sd == 0:
        print("✓ No straight double quotes remain (0.4 points)")
        score += 0.4
    else:
        print("✗ Straight double quotes still present (0 points)")

    # B) Straight single quotes gone
    if ss == 0:
        print("✓ No straight single quotes remain (0.4 points)")
        score += 0.4
    else:
        print("✗ Straight single quotes still present (0 points)")

    # C) At least one curly quote present
    if (cs + cd) > 0:
        print("✓ Curly quotes detected (0.2 points)")
        score += 0.2
    else:
        print("✗ No curly quotes detected (0 points)")

    # --- Final score ---
    final_score = min(score, max_score)
    print(f"Total score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score


# Execute verification when run as a script
if __name__ == "__main__":
    FILE_PATH = "/home/user/replace_straight_quotes_with_typographic_quotes_for_the_whole_document.pptx"
    verify_typographic_quotes(FILE_PATH)

