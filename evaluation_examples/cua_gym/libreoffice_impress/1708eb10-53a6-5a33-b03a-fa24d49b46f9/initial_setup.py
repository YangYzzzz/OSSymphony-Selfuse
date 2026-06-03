"""
Initial Setup: Create a 7-slide study review presentation with slides 3-5 empty
Task ID: impress_stu_064
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_064'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, alignment=PP_ALIGN.LEFT, color=None):
    """Helper to add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    add_text_box(slide1, Inches(1.5), Inches(1.5), Inches(10), Inches(2),
                 "Biology & Science Study Review", font_size=40, bold=True,
                 alignment=PP_ALIGN.CENTER, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text_box(slide1, Inches(2), Inches(4), Inches(9), Inches(1.5),
                 "Weekly Study Group Session - Spring 2025\nPrepared by: Emily Rodriguez",
                 font_size=22, alignment=PP_ALIGN.CENTER,
                 color=RGBColor(0xBD, 0xC3, 0xC7))

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide2, Inches(0.8), Inches(0.5), Inches(8), Inches(1),
                 "Session Agenda", font_size=32, bold=True,
                 color=RGBColor(0x2C, 0x3E, 0x50))

    agenda_items = [
        "1. Review key concepts from Chapter 7-9",
        "2. Interactive Quiz Section (Slides 3-5)",
        "3. Discussion of difficult topics",
        "4. Practice problem walkthrough",
        "5. Next week's reading assignments",
    ]
    y_pos = Inches(1.8)
    for item in agenda_items:
        add_text_box(slide2, Inches(1.2), y_pos, Inches(10), Inches(0.5),
                     item, font_size=20, color=RGBColor(0x34, 0x49, 0x5E))
        y_pos += Inches(0.7)

    # --- Slide 3: EMPTY (quiz placeholder) ---
    prs.slides.add_slide(prs.slide_layouts[5])

    # --- Slide 4: EMPTY (quiz placeholder) ---
    prs.slides.add_slide(prs.slide_layouts[5])

    # --- Slide 5: EMPTY (quiz placeholder) ---
    prs.slides.add_slide(prs.slide_layouts[5])

    # --- Slide 6: Key Takeaways ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide6, Inches(0.8), Inches(0.5), Inches(8), Inches(1),
                 "Key Takeaways", font_size=32, bold=True,
                 color=RGBColor(0x2C, 0x3E, 0x50))

    takeaways = [
        "Cellular biology fundamentals are the building blocks of life sciences",
        "Planetary science helps us understand our place in the solar system",
        "Chemistry vocabulary is essential for advanced coursework",
        "Practice quizzes improve retention by 40% compared to passive review",
    ]
    y_pos = Inches(1.8)
    for item in takeaways:
        add_text_box(slide6, Inches(1.2), y_pos, Inches(10), Inches(0.6),
                     f"\u2022 {item}", font_size=18, color=RGBColor(0x34, 0x49, 0x5E))
        y_pos += Inches(0.8)

    # --- Slide 7: Next Steps ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = slide7.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    add_text_box(slide7, Inches(2), Inches(2), Inches(9), Inches(1.5),
                 "Next Session: April 15, 2025", font_size=36, bold=True,
                 alignment=PP_ALIGN.CENTER, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text_box(slide7, Inches(2), Inches(4), Inches(9), Inches(1),
                 "Read Chapters 10-12 and complete practice worksheet",
                 font_size=22, alignment=PP_ALIGN.CENTER,
                 color=RGBColor(0xBD, 0xC3, 0xC7))
    add_text_box(slide7, Inches(2), Inches(5.5), Inches(9), Inches(0.8),
                 "Questions? Email: study-group@university.edu",
                 font_size=18, alignment=PP_ALIGN.CENTER,
                 color=RGBColor(0x95, 0xA5, 0xA6))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
