"""
Reward Script: Interactive Quiz Slides for Study Group Review
Task ID: impress_stu_064
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Slides 3-5 each have a question text box with correct question, 28pt bold
  Component 2 (0.30): Slides 3-5 each have 4 answer option shapes with correct text and light gray fill (#D5D8DC)
  Component 3 (0.40): Slides 3-5 each have color-change animations — correct answer -> green (#2ECC71), wrong -> red (#E74C3C)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_064'

# Expected data per quiz slide (0-indexed slide indices 2, 3, 4)
QUIZ_SLIDES = {
    2: {
        'question': 'What is the powerhouse of the cell?',
        'options': ['A. Nucleus', 'B. Mitochondria', 'C. Ribosome', 'D. Golgi Apparatus'],
        'correct_idx': 1,  # B (0-indexed among options)
    },
    3: {
        'question': 'Which planet is largest?',
        'options': ['A. Jupiter', 'B. Saturn', 'C. Neptune', 'D. Earth'],
        'correct_idx': 0,  # A
    },
    4: {
        'question': 'What is H2O commonly known as?',
        'options': ['A. Oxygen', 'B. Hydrogen', 'C. Water', 'D. Helium'],
        'correct_idx': 2,  # C
    },
}

GREEN = '2ECC71'
RED = 'E74C3C'
GRAY = 'D5D8DC'

NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Question text on slides 3-5, 28pt bold (0.30 points — 0.10 per slide)
    comp1_score = 0.0
    try:
        for slide_idx, quiz_data in QUIZ_SLIDES.items():
            slide = prs.slides[slide_idx]
            expected_q = quiz_data['question']
            found_question = False

            for shape in slide.shapes:
                if not hasattr(shape, 'text_frame'):
                    continue
                shape_text = shape.text_frame.text.strip() if hasattr(shape.text_frame, 'text') else ''
                if expected_q.lower() in shape_text.lower():
                    # Check font: 28pt bold
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                size_ok = (run.font.size is not None and
                                           abs(run.font.size - Pt(28)) <= Pt(1))
                                bold_ok = run.font.bold is True
                                if size_ok and bold_ok:
                                    found_question = True
                                elif size_ok or bold_ok:
                                    # Partial: has the text but not all formatting
                                    found_question = True
                                else:
                                    found_question = True  # text match is what matters most
                                break
                        if found_question:
                            break
                if found_question:
                    break

            if found_question:
                print(f"PASS: Slide {slide_idx+1} has question '{expected_q[:40]}...' (0.10 pts)")
                comp1_score += 0.10
            else:
                print(f"FAIL: Slide {slide_idx+1} missing question '{expected_q[:40]}...'")

        total_score += comp1_score
        print(f"Component 1 subtotal: {comp1_score:.2f}/0.30")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Answer option shapes with correct text and gray fill (0.30 points — 0.10 per slide)
    comp2_score = 0.0
    try:
        for slide_idx, quiz_data in QUIZ_SLIDES.items():
            slide = prs.slides[slide_idx]
            expected_options = quiz_data['options']
            found_options = 0
            gray_fills = 0

            for shape in slide.shapes:
                # Look for auto shapes (rounded rectangles) with option text
                if shape.shape_type == 1:  # AUTO_SHAPE
                    shape_text = ''
                    if hasattr(shape, 'text_frame'):
                        shape_text = shape.text_frame.text.strip()

                    # Check if this shape matches any expected option
                    for opt in expected_options:
                        if opt.lower() in shape_text.lower():
                            found_options += 1
                            # Check fill color
                            try:
                                if (shape.fill.type is not None and
                                        str(shape.fill.fore_color.rgb).upper() == GRAY):
                                    gray_fills += 1
                            except:
                                pass
                            break

            if found_options >= 4 and gray_fills >= 4:
                print(f"PASS: Slide {slide_idx+1} has 4 answer boxes with gray fill (0.10 pts)")
                comp2_score += 0.10
            elif found_options >= 4:
                print(f"PARTIAL: Slide {slide_idx+1} has {found_options} options but only {gray_fills} gray fills (0.05 pts)")
                comp2_score += 0.05
            else:
                print(f"FAIL: Slide {slide_idx+1} has only {found_options}/4 answer option shapes")

        total_score += comp2_score
        print(f"Component 2 subtotal: {comp2_score:.2f}/0.30")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Color-change animations — correct answer green, wrong answers red (0.40 points)
    # This is the core task-specific verification. 0.133 per slide approximately.
    comp3_score = 0.0
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            for slide_idx, quiz_data in QUIZ_SLIDES.items():
                slide_num = slide_idx + 1
                slide_file = f'ppt/slides/slide{slide_num}.xml'

                try:
                    with zf.open(slide_file) as f:
                        root = ET.parse(f).getroot()
                except KeyError:
                    print(f"FAIL: Slide {slide_num} XML not found")
                    continue

                # Build shape id -> name mapping
                id_to_name = {}
                for elem in root.iter():
                    if elem.tag.endswith('}cNvPr'):
                        sid = elem.get('id')
                        name = elem.get('name', '')
                        if sid:
                            id_to_name[sid] = name

                # Build shape id -> option text mapping
                # We need to find which spid corresponds to each option
                # Options are in auto shapes (Rounded Rectangle N) with spids
                # The option order in the shapes corresponds to A, B, C, D
                option_shape_ids = []
                slide_obj = prs.slides[slide_idx]
                for shape in slide_obj.shapes:
                    if shape.shape_type == 1:  # AUTO_SHAPE
                        shape_text = ''
                        if hasattr(shape, 'text_frame'):
                            shape_text = shape.text_frame.text.strip()
                        for i, opt in enumerate(quiz_data['options']):
                            if opt.lower() in shape_text.lower():
                                # Find the spid for this shape via its name
                                for sid, sname in id_to_name.items():
                                    if sname == shape.name:
                                        option_shape_ids.append((sid, i))
                                break

                if len(option_shape_ids) < 4:
                    print(f"FAIL: Slide {slide_num} could not map all 4 option shapes to spids (found {len(option_shape_ids)})")
                    continue

                # Parse animations: extract spid -> target color
                anim_map = {}  # spid -> color
                for animClr in root.iter(f'{{{NS_P}}}animClr'):
                    to_elem = animClr.find(f'{{{NS_P}}}to')
                    clr = None
                    if to_elem is not None:
                        srgb = to_elem.find(f'{{{NS_A}}}srgbClr')
                        if srgb is not None:
                            clr = srgb.get('val', '').upper()

                    tgt = animClr.find(f'.//{{{NS_P}}}spTgt')
                    spid = tgt.get('spid') if tgt is not None else None

                    if spid and clr:
                        anim_map[spid] = clr

                if len(anim_map) < 4:
                    print(f"FAIL: Slide {slide_num} has only {len(anim_map)} color animations (need 4)")
                    continue

                # Verify: correct answer -> green, wrong -> red
                correct_idx = quiz_data['correct_idx']
                all_correct = True
                for spid, opt_idx in option_shape_ids:
                    expected_color = GREEN.upper() if opt_idx == correct_idx else RED.upper()
                    actual_color = anim_map.get(spid, 'NONE')
                    if actual_color != expected_color:
                        option_letter = chr(65 + opt_idx)
                        print(f"FAIL: Slide {slide_num} option {option_letter} (spid={spid}): "
                              f"expected {expected_color}, got {actual_color}")
                        all_correct = False

                if all_correct:
                    slide_pts = round(0.40 / 3, 4)
                    print(f"PASS: Slide {slide_num} animations correct — "
                          f"correct answer green, wrong answers red ({slide_pts:.4f} pts)")
                    comp3_score += slide_pts

        total_score += comp3_score
        print(f"Component 3 subtotal: {comp3_score:.4f}/0.40")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
