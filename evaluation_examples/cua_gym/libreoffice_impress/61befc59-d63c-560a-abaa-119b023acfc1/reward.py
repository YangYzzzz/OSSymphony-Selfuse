"""
Reward Script: Add stacked bar chart on slide 4 with assignment scores
Task ID: impress_teach_049
Domain: libreoffice_impress
Scoring:
  Component 1: Chart exists on slide 4 and is BAR_STACKED type (0.25)
  Component 2: Chart title is 'Score Breakdown by Student' (0.15)
  Component 3: Categories are Student A, Student B, Student C (0.15)
  Component 4: Series names are Homework, Midterm, Final (0.15)
  Component 5: Series data values match expected values (0.30)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_049'

# Expected data
EXPECTED_CATEGORIES = ['Student A', 'Student B', 'Student C']
EXPECTED_SERIES = {
    'Homework': [80.0, 90.0, 70.0],
    'Midterm': [75.0, 88.0, 65.0],
    'Final': [85.0, 92.0, 78.0],
}
EXPECTED_TITLE = 'Score Breakdown by Student'


def persist_app_state(domain):
    """Save any unsaved GUI edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_series_name(series):
    """Extract series name from XML element."""
    ns = {'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'}
    tx = series._element.find('.//c:tx//c:v', ns)
    if tx is not None:
        return tx.text
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation
    from pptx.enum.chart import XL_CHART_TYPE

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # slide 4 (0-indexed)

    # Find chart shape on slide 4
    chart_shape = None
    for shape in slide.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Chart exists on slide 4 and is BAR_STACKED type (0.25 points)
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            if chart.chart_type == XL_CHART_TYPE.BAR_STACKED:
                print(f"PASS: Component 1 - Stacked bar chart found on slide 4 (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 - Chart found but type is {chart.chart_type}, expected BAR_STACKED")
        else:
            print(f"FAIL: Component 1 - No chart found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # If no chart, remaining components cannot pass
    if chart_shape is None:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    chart = chart_shape.chart

    # Component 2: Chart title is 'Score Breakdown by Student' (0.15 points)
    try:
        if chart.has_title:
            actual_title = chart.chart_title.text_frame.text.strip()
            if actual_title == EXPECTED_TITLE:
                print(f"PASS: Component 2 - Chart title matches: '{actual_title}' (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 2 - Chart title is '{actual_title}', expected '{EXPECTED_TITLE}'")
        else:
            print(f"FAIL: Component 2 - Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Categories are Student A, Student B, Student C (0.15 points)
    try:
        plot = chart.plots[0]
        actual_cats = [str(c) for c in plot.categories]
        if actual_cats == EXPECTED_CATEGORIES:
            print(f"PASS: Component 3 - Categories match: {actual_cats} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 3 - Categories are {actual_cats}, expected {EXPECTED_CATEGORIES}")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Series names are Homework, Midterm, Final (0.15 points)
    try:
        actual_series_names = []
        for series in chart.series:
            name = get_series_name(series)
            actual_series_names.append(name)
        expected_names = list(EXPECTED_SERIES.keys())
        if actual_series_names == expected_names:
            print(f"PASS: Component 4 - Series names match: {actual_series_names} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 - Series names are {actual_series_names}, expected {expected_names}")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Series data values match expected (0.30 points)
    # Award partial credit: 0.10 per correct series
    try:
        series_score = 0.0
        expected_names = list(EXPECTED_SERIES.keys())
        for i, series in enumerate(chart.series):
            s_name = get_series_name(series)
            actual_vals = list(series.values)
            if s_name in EXPECTED_SERIES:
                expected_vals = EXPECTED_SERIES[s_name]
                if actual_vals == expected_vals:
                    print(f"PASS: Component 5.{i+1} - Series '{s_name}' values match: {actual_vals} (0.10 pts)")
                    series_score += 0.10
                else:
                    print(f"FAIL: Component 5.{i+1} - Series '{s_name}' values are {actual_vals}, expected {expected_vals}")
            elif i < len(expected_names):
                # Try matching by index
                exp_name = expected_names[i]
                expected_vals = EXPECTED_SERIES[exp_name]
                if actual_vals == expected_vals:
                    print(f"PASS: Component 5.{i+1} - Series index {i} values match (name '{s_name}' vs expected '{exp_name}'): {actual_vals} (0.10 pts)")
                    series_score += 0.10
                else:
                    print(f"FAIL: Component 5.{i+1} - Series index {i} values are {actual_vals}, expected {expected_vals}")
            else:
                print(f"FAIL: Component 5.{i+1} - Extra series '{s_name}' not expected")

        if series_score > 0:
            total_score += series_score
            print(f"Component 5 subtotal: {series_score}/0.30")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
