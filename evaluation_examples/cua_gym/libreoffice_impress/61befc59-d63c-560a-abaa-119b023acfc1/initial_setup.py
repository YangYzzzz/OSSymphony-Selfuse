"""
Initial Setup: Add a stacked bar chart on slide 4
Task ID: impress_teach_049
Domain: libreoffice_impress

Creates a 6-slide presentation about a gradebook.
Slide 4 has the title 'Detailed Scores' but NO chart (task requires adding one).
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_049'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_body(slide, title_text, body_lines):
    """Helper to set title and add body text to a slide."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    # Find body placeholder (index 1)
    body_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            body_ph = ph
            break
    if body_ph and body_lines:
        tf = body_ph.text_frame
        tf.clear()
        for i, line in enumerate(body_lines):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                p = tf.add_paragraph()
                p.text = line


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Gradebook Presentation"
    slide1.placeholders[1].text = "Fall 2025 - Introduction to Computer Science\nPrepared by Prof. Elena Rodriguez"

    # --- Slide 2: Course Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide2, "Course Overview", [
        "Course: CS 101 - Introduction to Computer Science",
        "Semester: Fall 2025",
        "Total Enrolled: 3 Students",
        "Assessment Components: Homework (30%), Midterm (30%), Final (40%)",
        "Grading Scale: A (90-100), B (80-89), C (70-79), D (60-69), F (<60)",
    ])

    # --- Slide 3: Student Roster ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide3, "Student Roster", [
        "Student A - Enrolled since Aug 2025",
        "Student B - Enrolled since Aug 2025",
        "Student C - Enrolled since Aug 2025",
    ])

    # --- Slide 4: Detailed Scores (NO CHART - task requires adding one) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide4, "Detailed Scores", [
        "Below is a summary of individual scores across all assessments.",
        "",
        "Student A: Homework 80, Midterm 75, Final 85",
        "Student B: Homework 90, Midterm 88, Final 92",
        "Student C: Homework 70, Midterm 65, Final 78",
    ])

    # --- Slide 5: Grade Summary ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide5, "Grade Summary", [
        "Student A - Total: 240 - Average: 80.0 - Grade: B",
        "Student B - Total: 270 - Average: 90.0 - Grade: A",
        "Student C - Total: 213 - Average: 71.0 - Grade: C",
        "",
        "Class Average: 80.3",
    ])

    # --- Slide 6: Thank You ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[0])
    slide6.shapes.title.text = "Thank You"
    slide6.placeholders[1].text = "Questions? Contact: e.rodriguez@university.edu"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
