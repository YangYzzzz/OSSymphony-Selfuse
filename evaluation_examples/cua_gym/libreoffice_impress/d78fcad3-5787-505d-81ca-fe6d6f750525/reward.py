"""
FINAL REWARD SCRIPT - SUCCESS
Task: Add a 3×4 table above paragraph 2 with a header row.
Generated: 2025-10-17 07:32:25
Status: success
Model: azure-o3
Total Steps: 7
"""

import os
import re
from pptx import Presentation


def verify_task(file_path: str) -> float:
    """Verify that the presentation contains a 3×4 table positioned
    above the shape containing *Paragraph 2*, and that the table is
    flagged with a header (first) row.

    Scoring (progressive):
        • 0.20 – Found a text shape whose content includes "Paragraph 2"
        • 0.40 – Found a 3×4 table on the same slide
        • 0.20 – Table’s first_row property == True  (header row)
        • 0.20 – Table’s top is above the Paragraph 2 text shape’s top
    Returns a float in [0.0, 1.0].
    """

    max_score = 1.0
    score = 0.0

    # Weights for each verification step
    W_PARAGRAPH = 0.20
    W_TABLE_DIM = 0.40
    W_HEADER    = 0.20
    W_POSITION  = 0.20

    # ---------- 1) Load presentation (no points just for loading) ----------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0
    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded presentation with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Unable to open presentation: {e}")
        return 0.0

    # ---------- 2) Locate shape containing "Paragraph 2" ----------
    target_slide = None
    paragraph_shape = None
    for s_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            combined_text = "\n".join(
                para.text for para in shape.text_frame.paragraphs if para.text
            )
            if re.search(r"paragraph\s*2", combined_text, re.IGNORECASE):
                target_slide = slide
                paragraph_shape = shape
                score += W_PARAGRAPH
                print(f"✓ Found 'Paragraph 2' text on slide {s_idx}")
                break
        if target_slide:
            break

    if not target_slide:
        print("✗ 'Paragraph 2' text not found in any slide")
        print(f"REWARD: {score}")
        return score  # Cannot continue without this reference

    # Y-coordinate of Paragraph 2 text box
    para_top = paragraph_shape.top

    # ---------- 3) Locate a 3×4 table on the same slide ----------
    table_shape = None
    for shape in target_slide.shapes:
        if not shape.has_table:
            continue
        tbl = shape.table
        if len(tbl.rows) == 3 and len(tbl.columns) == 4:
            table_shape = shape
            break

    if not table_shape:
        print("✗ No 3×4 table found on the slide containing 'Paragraph 2'")
        print(f"REWARD: {score}")
        return score

    print("✓ Found table with required 3×4 dimensions")
    score += W_TABLE_DIM

    # ---------- 4) Verify header row flag ----------
    has_header = getattr(table_shape.table, "first_row", False)
    if has_header:
        print("✓ Table 'first_row' property is set → header row detected")
        score += W_HEADER
    else:
        print("✗ Table does not have 'first_row' header flag set")

    # ---------- 5) Verify positional relationship ----------
    if table_shape.top < para_top:
        print("✓ Table is positioned above the 'Paragraph 2' text shape")
        score += W_POSITION
    else:
        print("✗ Table is not positioned above the 'Paragraph 2' text shape")

    # Cap to 1.0 and report
    final_score = min(score, max_score)
    print(f"REWARD: {final_score}")
    return final_score


# ---------------------------------------------------------------------------
# Execute verification when run as a standalone script
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    verify_task("/home/user/add_a_34_table_above_paragraph_2_with_a_header_row.pptx")
