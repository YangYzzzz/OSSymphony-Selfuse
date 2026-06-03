"""
Initial Setup: Create a Summary_Deck presentation with 8 slides.
Slide 7 has an arrow shape in top-left and 'Key Conclusion' text in bottom-right.
No animations.
Task ID: impress_ma_087
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_087'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q4 2025 Strategic Summary"
    slide1.placeholders[1].text = "Prepared by the Corporate Strategy Division"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txb = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.text = "Agenda"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    agenda_items = [
        "1. Revenue Performance Overview",
        "2. Market Expansion Analysis",
        "3. Product Launch Metrics",
        "4. Customer Retention Trends",
        "5. Operational Efficiency Gains",
        "6. Key Conclusions & Next Steps",
    ]
    agenda_box = slide2.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10), Inches(4.5))
    atf = agenda_box.text_frame
    atf.word_wrap = True
    for i, item in enumerate(agenda_items):
        para = atf.paragraphs[0] if i == 0 else atf.add_paragraph()
        para.text = item
        para.space_after = Pt(12)
        r = para.runs[0]
        r.font.size = Pt(22)
        r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 3: Revenue Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    title3 = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    title3.text_frame.paragraphs[0].text = "Revenue Performance"
    title3.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    title3.text_frame.paragraphs[0].runs[0].font.bold = True

    tbl_shape = slide3.shapes.add_table(5, 4, Inches(1.5), Inches(2), Inches(10), Inches(3.5))
    tbl = tbl_shape.table
    headers = ["Region", "Q3 Revenue ($M)", "Q4 Revenue ($M)", "Growth (%)"]
    data = [
        ["North America", "$142.3", "$158.7", "+11.5%"],
        ["Europe", "$98.1", "$107.4", "+9.5%"],
        ["Asia Pacific", "$67.5", "$79.2", "+17.3%"],
        ["Latin America", "$23.8", "$28.1", "+18.1%"],
    ]
    for c, h in enumerate(headers):
        tbl.cell(0, c).text = h
        for r in tbl.cell(0, c).text_frame.paragraphs[0].runs:
            r.font.bold = True
            r.font.size = Pt(14)
    for ri, row in enumerate(data, 1):
        for ci, val in enumerate(row):
            tbl.cell(ri, ci).text = val

    # --- Slide 4: Market Expansion ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    title4 = slide4.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    title4.text_frame.paragraphs[0].text = "Market Expansion Analysis"
    title4.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    title4.text_frame.paragraphs[0].runs[0].font.bold = True

    bullets4 = slide4.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10), Inches(4))
    btf4 = bullets4.text_frame
    btf4.word_wrap = True
    items4 = [
        "Entered 3 new markets in Southeast Asia (Vietnam, Philippines, Thailand)",
        "B2B partnerships grew by 24% compared to Q3",
        "Enterprise client acquisition increased from 47 to 63 accounts",
        "New verticals: Healthcare logistics and EdTech platforms",
    ]
    for i, item in enumerate(items4):
        para = btf4.paragraphs[0] if i == 0 else btf4.add_paragraph()
        para.text = item
        para.space_after = Pt(10)
        r = para.runs[0]
        r.font.size = Pt(20)

    # --- Slide 5: Product Launch ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    title5 = slide5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    title5.text_frame.paragraphs[0].text = "Product Launch Metrics"
    title5.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    title5.text_frame.paragraphs[0].runs[0].font.bold = True

    tbl5_shape = slide5.shapes.add_table(4, 3, Inches(2), Inches(2.2), Inches(9), Inches(3))
    tbl5 = tbl5_shape.table
    h5 = ["Product", "Launch Date", "Adoption Rate"]
    d5 = [
        ["CloudSync Pro", "Oct 1, 2025", "34% (target 30%)"],
        ["DataVault Enterprise", "Nov 15, 2025", "21% (target 25%)"],
        ["InsightFlow Analytics", "Dec 3, 2025", "28% (target 20%)"],
    ]
    for c, h in enumerate(h5):
        tbl5.cell(0, c).text = h
        for r in tbl5.cell(0, c).text_frame.paragraphs[0].runs:
            r.font.bold = True
    for ri, row in enumerate(d5, 1):
        for ci, val in enumerate(row):
            tbl5.cell(ri, ci).text = val

    # --- Slide 6: Customer Retention ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    title6 = slide6.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    title6.text_frame.paragraphs[0].text = "Customer Retention Trends"
    title6.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    title6.text_frame.paragraphs[0].runs[0].font.bold = True

    bullets6 = slide6.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10), Inches(4))
    btf6 = bullets6.text_frame
    btf6.word_wrap = True
    items6 = [
        "Overall retention rate: 91.2% (up from 88.7% in Q3)",
        "Net Promoter Score improved to 72 from 65",
        "Churn rate decreased to 4.3% for enterprise tier",
        "Customer lifetime value grew 15% YoY to $18,400",
        "Top churn reasons: pricing (32%), feature gaps (28%), competition (21%)",
    ]
    for i, item in enumerate(items6):
        para = btf6.paragraphs[0] if i == 0 else btf6.add_paragraph()
        para.text = item
        para.space_after = Pt(10)
        r = para.runs[0]
        r.font.size = Pt(20)

    # --- Slide 7: Key Conclusion (THE TASK SLIDE) ---
    # Arrow shape in top-left, "Key Conclusion" text box in bottom-right
    # NO animations
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    title7 = slide7.shapes.add_textbox(Inches(1), Inches(0.3), Inches(11), Inches(0.8))
    title7.text_frame.paragraphs[0].text = "Strategic Direction"
    title7.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    title7.text_frame.paragraphs[0].runs[0].font.bold = True
    title7.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # Arrow shape in TOP-LEFT corner (this is the shape to be animated)
    arrow = slide7.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(0.8), Inches(1.5),  # top-left position
        Inches(1.5), Inches(0.8),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0xE7, 0x4C, 0x3C)
    arrow.name = "PointerArrow"

    # "Key Conclusion" text box in BOTTOM-RIGHT
    conclusion_box = slide7.shapes.add_textbox(
        Inches(7.5), Inches(4.5),
        Inches(5), Inches(2.5),
    )
    ctf = conclusion_box.text_frame
    ctf.word_wrap = True
    p_title = ctf.paragraphs[0]
    p_title.text = "Key Conclusion"
    p_title.alignment = PP_ALIGN.LEFT
    r_title = p_title.runs[0]
    r_title.font.size = Pt(28)
    r_title.font.bold = True
    r_title.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    p_body = ctf.add_paragraph()
    p_body.text = (
        "Strong cross-regional growth and successful product launches "
        "position the company for sustained expansion in 2026. "
        "Focus areas: deepen Asia Pacific presence and accelerate "
        "enterprise adoption of DataVault."
    )
    r_body = p_body.runs[0]
    r_body.font.size = Pt(18)
    r_body.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    conclusion_box.name = "ConclusionText"

    # --- Slide 8: Next Steps ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    title8 = slide8.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    title8.text_frame.paragraphs[0].text = "Next Steps & Action Items"
    title8.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    title8.text_frame.paragraphs[0].runs[0].font.bold = True

    bullets8 = slide8.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10), Inches(4))
    btf8 = bullets8.text_frame
    btf8.word_wrap = True
    items8 = [
        "Finalize Q1 2026 hiring plan for APAC expansion (Due: Jan 15)",
        "Launch DataVault Enterprise 2.0 beta program (Due: Feb 1)",
        "Complete pricing review with finance team (Due: Jan 30)",
        "Schedule quarterly business reviews with top 20 enterprise clients",
        "Submit board presentation with updated 3-year projections",
    ]
    for i, item in enumerate(items8):
        para = btf8.paragraphs[0] if i == 0 else btf8.add_paragraph()
        para.text = item
        para.space_after = Pt(10)
        r = para.runs[0]
        r.font.size = Pt(20)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
