"""
Initial Setup: 5-slide executive briefing presentation
Task ID: osworld_impress_underline_darkred_table_012
Domain: libreoffice_impress

Creates a 5-slide executive briefing .pptx where slide 4 has a content
textbox with 5 bullet items in black (no underline, no #5C0000 color).
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_underline_darkred_table_012'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_content_slide(prs, title_text, content_lines, layout_idx=1):
    """Add a slide with a title placeholder and a content placeholder."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    # Set title
    title_ph = slide.placeholders[0]
    title_ph.text = title_text
    tf_title = title_ph.text_frame
    for para in tf_title.paragraphs:
        for run in para.runs:
            run.font.bold = True
            run.font.size = Pt(32)
            run.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)

    # Set content
    content_ph = slide.placeholders[1]
    content_ph.text = ""
    tf = content_ph.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.text = line
        para.level = 0
        for run in para.runs:
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
            run.font.bold = False
            run.font.italic = False
            run.font.underline = False

    return slide


def create_initial():
    prs = Presentation()
    # Use default slide dimensions (widescreen 10" x 7.5")

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.placeholders[0].text = "Q3 2025 Executive Briefing"
    slide1.placeholders[1].text = "Strategic Initiatives & Performance Review"
    for run in slide1.placeholders[0].text_frame.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)
        run.font.size = Pt(40)
        run.font.bold = True
    for run in slide1.placeholders[1].text_frame.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0x44, 0x47, 0x2C)
        run.font.size = Pt(24)

    # --- Slide 2: Financial Summary ---
    add_title_content_slide(
        prs,
        title_text="Financial Summary",
        content_lines=[
            "Total Revenue: $4.7M (up 12% YoY)",
            "Operating Expenses: $3.1M (down 4% from Q2)",
            "Net Profit Margin: 18.2%",
            "EBITDA: $1.4M, exceeding guidance by $120K",
            "Cash on Hand: $6.8M — runway of 22 months",
        ]
    )

    # --- Slide 3: Operational Highlights ---
    add_title_content_slide(
        prs,
        title_text="Operational Highlights",
        content_lines=[
            "Launched 3 new product features in July & August",
            "Customer satisfaction score rose to 87% (from 81%)",
            "Support ticket resolution time reduced by 30%",
            "Engineering team expanded to 42 full-time employees",
            "Partnership agreements signed with DataBridge and NovaCorp",
        ]
    )

    # --- Slide 4: Key Risks & Mitigation ---
    # This is the target slide — 5 bullet items in black, no underline, no dark red
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title_ph4 = slide4.placeholders[0]
    title_ph4.text = "Key Risks & Mitigation"
    for para in title_ph4.text_frame.paragraphs:
        for run in para.runs:
            run.font.bold = True
            run.font.size = Pt(32)
            run.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)

    content_ph4 = slide4.placeholders[1]
    content_ph4.text = ""
    tf4 = content_ph4.text_frame
    tf4.word_wrap = True

    bullet_items = [
        "Supply chain disruptions may delay hardware deliveries by 2–4 weeks",
        "Regulatory compliance review scheduled for October 2025",
        "Competitive pressure increasing from three new market entrants",
        "Talent retention risk in engineering — proactive raise cycle planned",
        "Cloud infrastructure costs projected to rise 15% in Q4 2025",
    ]

    for i, line in enumerate(bullet_items):
        if i == 0:
            para = tf4.paragraphs[0]
        else:
            para = tf4.add_paragraph()
        para.text = line
        para.level = 0
        for run in para.runs:
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black — NOT dark red
            run.font.bold = False
            run.font.italic = False
            run.font.underline = False  # NO underline

    # --- Slide 5: Next Steps ---
    add_title_content_slide(
        prs,
        title_text="Next Steps",
        content_lines=[
            "Finalize Q4 budget allocations by September 15",
            "Complete security audit and submit compliance report",
            "Kick off Phase 2 of the platform migration project",
            "Hire 5 additional engineers and 2 product managers",
            "Schedule board review meeting for October 2025",
        ]
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
