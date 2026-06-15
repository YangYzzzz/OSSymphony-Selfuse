"""
Reward Script: Underline all bullet points in content textbox on slide 4 and change color to dark maroon (#5C0000)
Task ID: osworld_impress_underline_darkred_table_012
Domain: libreoffice_impress
Scoring:
  Component 1: All 5 bullet items in content textbox on slide 4 have underline=True (0.5 points)
  Component 2: All 5 bullet items in content textbox on slide 4 have color=#5C0000 (0.5 points)
  Total: 1.0
"""

import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_underline_darkred_table_012'
TARGET_COLOR = '5C0000'  # dark maroon #5C0000
EXPECTED_BULLETS = 5
SLIDE_INDEX = 3  # Slide 4 (0-indexed)


def persist_app_state():
    """Send ctrl+s to persist any unsaved edits from LibreOffice."""
    try:
        import time
        os.environ["DISPLAY"] = ":0"
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0

    The task requires:
    - All 5 bullet items in Content Placeholder 2 on slide 4 have underline=True
    - All 5 bullet items have their text color changed to #5C0000 (dark maroon)
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Verify presentation has at least 4 slides
    if len(prs.slides) < 4:
        print(f"CRITICAL: Expected at least 4 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[SLIDE_INDEX]

    # Find the content textbox (Content Placeholder) on slide 4
    content_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == 'Content Placeholder 2':
            content_shape = shape
            break

    # Fallback: use any non-title placeholder that has bullet-like content
    if content_shape is None:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                # Skip title shapes (usually shorter text, no bullets)
                paras = shape.text_frame.paragraphs
                if len(paras) >= 3:
                    content_shape = shape
                    break

    if content_shape is None:
        print("CRITICAL: Content textbox with bullet points not found on slide 4")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Found content shape: {content_shape.name}")
    paragraphs = content_shape.text_frame.paragraphs

    # Filter to non-empty paragraphs (bullet items)
    bullet_paras = [p for p in paragraphs if p.text.strip()]
    print(f"INFO: Found {len(bullet_paras)} non-empty paragraphs (bullet items)")

    # ---- Component 1: All bullet items have underline=True (0.5 points) ----
    try:
        underline_count = 0
        underline_total = 0
        for para in bullet_paras:
            runs = [r for r in para.runs if (r.text or "").strip()]
            if not runs:
                # If no runs but text present, count this bullet
                underline_total += 1
                print(f"WARN: Para '{para.text[:40]}...' has no runs — cannot check underline")
                continue
            for run in runs:
                underline_total += 1
                if run.font.underline is True:
                    underline_count += 1
                else:
                    print(f"FAIL detail: Run '{run.text[:40]}' underline={run.font.underline} (expected True)")

        if underline_total > 0 and underline_count == underline_total:
            print(f"PASS: Component 1 — All {underline_count}/{underline_total} runs have underline=True (0.5 pts)")
            total_score += 0.5
        elif underline_count > 0:
            partial = round(0.5 * (underline_count / underline_total), 3)
            print(f"PARTIAL: Component 1 — {underline_count}/{underline_total} runs have underline=True ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — No runs have underline=True (0/{underline_total})")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # ---- Component 2: All bullet items have color=#5C0000 (0.5 points) ----
    try:
        color_count = 0
        color_total = 0
        for para in bullet_paras:
            runs = [r for r in para.runs if (r.text or "").strip()]
            if not runs:
                color_total += 1
                print(f"WARN: Para '{para.text[:40]}...' has no runs — cannot check color")
                continue
            for run in runs:
                color_total += 1
                try:
                    color_type = run.font.color.type
                    if color_type is not None:
                        rgb = str(run.font.color.rgb).upper()
                        if rgb == TARGET_COLOR.upper():
                            color_count += 1
                        else:
                            print(f"FAIL detail: Run '{run.text[:40]}' color={rgb} (expected {TARGET_COLOR})")
                    else:
                        print(f"FAIL detail: Run '{run.text[:40]}' color type is None/inherited (expected #{TARGET_COLOR})")
                except Exception as ce:
                    print(f"WARN: Could not read color for run '{run.text[:40]}': {ce}")

        if color_total > 0 and color_count == color_total:
            print(f"PASS: Component 2 — All {color_count}/{color_total} runs have color=#{TARGET_COLOR} (0.5 pts)")
            total_score += 0.5
        elif color_count > 0:
            partial = round(0.5 * (color_count / color_total), 3)
            print(f"PARTIAL: Component 2 — {color_count}/{color_total} runs have color=#{TARGET_COLOR} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No runs have color=#{TARGET_COLOR} (0/{color_total})")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(round(total_score, 4), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
