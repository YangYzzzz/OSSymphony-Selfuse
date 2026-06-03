"""
Initial Setup: Create a 10-slide portfolio presentation and open it in LibreOffice Impress
Task ID: impress_el_016
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_el_016'
PPTX_PATH = f'{WORKDIR}/{TASK_ID}.pptx'
ODP_PATH = f'{WORKDIR}/Portfolio.odp'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Creative Portfolio"
    slide1.placeholders[1].text = "Elena Vasquez | Visual Designer & Illustrator"
    bg1 = slide1.background.fill
    bg1.solid()
    bg1.fore_color.rgb = RGBColor(0x1B, 0x2A, 0x4A)
    for para in slide1.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.size = Pt(40)
    for para in slide1.placeholders[1].text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            run.font.size = Pt(20)

    # --- Slide 2: About Me ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(1))
    p = txBox.text_frame.paragraphs[0]
    p.text = "About Me"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    txBox2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8), Inches(4))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = (
        "With over 8 years of experience in visual design, I specialize in brand identity, "
        "UI/UX design, and digital illustration. My work bridges the gap between aesthetics "
        "and functionality, creating designs that resonate with audiences and drive engagement."
    )
    for run in tf2.paragraphs[0].runs:
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # --- Slide 3: Design Philosophy ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    bg3 = slide3.background.fill
    bg3.solid()
    bg3.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    txBox = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    p = txBox.text_frame.paragraphs[0]
    p.text = "Design Philosophy"
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    items = [
        "Simplicity is the ultimate sophistication",
        "Every pixel should serve a purpose",
        "User experience drives every design decision",
        "Consistency builds trust and recognition",
    ]
    txBox2 = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tf = txBox2.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = f"   {item}"
        para.level = 0
        for run in para.runs:
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # --- Slide 4: Brand Identity Projects ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide4.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(1))
    p = txBox.text_frame.paragraphs[0]
    p.text = "Brand Identity Projects"
    run = p.runs[0]
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)
    # Table of projects
    table_shape = slide4.shapes.add_table(5, 3, Inches(0.8), Inches(2), Inches(8), Inches(3.5))
    table = table_shape.table
    headers = ["Client", "Project", "Year"]
    data = [
        ["Horizon Tech", "Full brand redesign", "2024"],
        ["Bloom Wellness", "Logo & packaging", "2023"],
        ["Metro Transit Authority", "Wayfinding system", "2023"],
        ["Solara Energy", "Corporate identity", "2022"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 5: UI/UX Showcase ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    bg5 = slide5.background.fill
    bg5.solid()
    bg5.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)
    txBox = slide5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    p = txBox.text_frame.paragraphs[0]
    p.text = "UI/UX Showcase"
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    txBox2 = slide5.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tf = txBox2.text_frame
    tf.word_wrap = True
    projects = [
        "FinTrack Mobile App - Personal finance dashboard with real-time data visualization",
        "CookBook Pro - Recipe management platform with AI-powered meal planning",
        "TravelMate - Trip planning app with collaborative itinerary features",
    ]
    for i, proj in enumerate(projects):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = proj
        for run in para.runs:
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)

    # --- Slide 6: Digital Illustrations ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide6.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    p = txBox.text_frame.paragraphs[0]
    p.text = "Digital Illustrations"
    run = p.runs[0]
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xE7, 0x4C, 0x3C)
    txBox2 = slide6.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tf = txBox2.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = (
        "Exploring themes of nature, urban landscapes, and human connection through digital media. "
        "Published in Artistry Magazine (2024), featured in the Digital Arts International Exhibition."
    )
    for run in tf.paragraphs[0].runs:
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # --- Slide 7: Awards & Recognition ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    bg7 = slide7.background.fill
    bg7.solid()
    bg7.fore_color.rgb = RGBColor(0xFD, 0xF2, 0xE9)
    txBox = slide7.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    p = txBox.text_frame.paragraphs[0]
    p.text = "Awards & Recognition"
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xD4, 0x7A, 0x2A)
    awards = [
        "Red Dot Design Award - Brand Identity (2024)",
        "Webby Award - Best Visual Design (2023)",
        "AIGA 50 Books | 50 Covers Selection (2023)",
        "Communication Arts Award of Excellence (2022)",
        "Type Directors Club Certificate of Typographic Excellence (2022)",
    ]
    txBox2 = slide7.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tf = txBox2.text_frame
    tf.word_wrap = True
    for i, award in enumerate(awards):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = award
        for run in para.runs:
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 8: Client Testimonials ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide8.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    p = txBox.text_frame.paragraphs[0]
    p.text = "Client Testimonials"
    run = p.runs[0]
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    quotes = [
        ('"Elena transformed our brand from forgettable to unforgettable. Her attention to detail is extraordinary."',
         "- James Park, CEO, Horizon Tech"),
        ('"The UI designs exceeded our expectations. User engagement increased by 47% after the redesign."',
         "- Maria Santos, Product Lead, FinTrack"),
    ]
    y_pos = Inches(2)
    for quote_text, attribution in quotes:
        txBox3 = slide8.shapes.add_textbox(Inches(1), y_pos, Inches(8), Inches(2))
        tf = txBox3.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = quote_text
        for run in tf.paragraphs[0].runs:
            run.font.size = Pt(16)
            run.font.italic = True
            run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
        p2 = tf.add_paragraph()
        p2.text = attribution
        p2.alignment = PP_ALIGN.RIGHT
        for run in p2.runs:
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        y_pos += Inches(2.5)

    # --- Slide 9: Skills & Tools ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    bg9 = slide9.background.fill
    bg9.solid()
    bg9.fore_color.rgb = RGBColor(0xEC, 0xF0, 0xF1)
    txBox = slide9.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    p = txBox.text_frame.paragraphs[0]
    p.text = "Skills & Tools"
    run = p.runs[0]
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    skills_data = [
        ["Design", "Adobe Creative Suite, Figma, Sketch, InVision"],
        ["Illustration", "Procreate, Adobe Illustrator, Wacom tablets"],
        ["Development", "HTML/CSS, React basics, Webflow, Framer"],
        ["Other", "Motion graphics, Photography, Typography"],
    ]
    table_shape = slide9.shapes.add_table(len(skills_data) + 1, 2, Inches(1), Inches(2), Inches(8), Inches(3))
    table = table_shape.table
    for c, h in enumerate(["Category", "Proficiencies"]):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, (cat, skills) in enumerate(skills_data, 1):
        table.cell(r, 0).text = cat
        table.cell(r, 1).text = skills

    # --- Slide 10: Contact ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    bg10 = slide10.background.fill
    bg10.solid()
    bg10.fore_color.rgb = RGBColor(0x1B, 0x2A, 0x4A)
    txBox = slide10.shapes.add_textbox(Inches(2), Inches(1.5), Inches(6), Inches(1.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "Let's Work Together"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    contact_info = [
        "elena.vasquez@designstudio.com",
        "www.elenavasquez.design",
        "+1 (415) 555-0187",
    ]
    txBox2 = slide10.shapes.add_textbox(Inches(2), Inches(3.5), Inches(6), Inches(3))
    tf = txBox2.text_frame
    tf.word_wrap = True
    for i, info in enumerate(contact_info):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = info
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)

    # Save as .pptx first
    prs.save(PPTX_PATH)
    print(f'Created pptx: {PPTX_PATH}')

    # Convert to ODP using LibreOffice CLI
    result = subprocess.run(
        ['libreoffice', '--headless', '--convert-to', 'odp', '--outdir', WORKDIR, PPTX_PATH],
        capture_output=True, text=True, timeout=30,
        env={**os.environ, "HOME": WORKDIR}
    )
    print(f'Convert stdout: {result.stdout}')
    print(f'Convert stderr: {result.stderr}')

    # Rename if needed (libreoffice outputs impress_el_016.odp, we want Portfolio.odp)
    converted_odp = f'{WORKDIR}/{TASK_ID}.odp'
    if os.path.exists(converted_odp):
        os.rename(converted_odp, ODP_PATH)
        print(f'Renamed to: {ODP_PATH}')
    elif os.path.exists(ODP_PATH):
        print(f'ODP already exists: {ODP_PATH}')
    else:
        print('WARNING: ODP conversion may have failed, checking...')
        print(subprocess.run(['ls', '-la', WORKDIR], capture_output=True, text=True).stdout)

    # Open in LibreOffice Impress for the GUI agent
    launch_gui(f'libreoffice --impress "{ODP_PATH}"', delay_sec=3.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
