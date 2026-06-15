"""
FINAL REWARD SCRIPT - SUCCESS
Task: In LibreOffice Impress, I’m polishing slide 115 and want the title text to pop a bit more—specifically, I need its underline switched from the default single style to a double line. How do I apply that change?
Generated: 2025-09-10 16:12:14
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
from pptx import Presentation
from pptx.enum.text import MSO_UNDERLINE
from pptx.enum.shapes import PP_PLACEHOLDER


def verify_double_underline_title(file_path: str, slide_number: int = 115) -> float:
    """Verify that the title on the specified slide uses DOUBLE underline
    (and no SINGLE underline remains). Returns a progressive score 0–1.0.
    """
    print(f"Verifying double underline for slide {slide_number} title in file: {file_path}")

    max_score = 1.0
    score = 0.0  # progressive scoring accumulator

    # 1) File existence & loadability (no points – prerequisite)
    if not os.path.exists(file_path):
        print("✗ File not found")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        return 0.0

    # 2) Slide count sufficiency (0.2 points)
    if len(prs.slides) < slide_number:
        print(f"✗ Presentation has only {len(prs.slides)} slides, expected ≥ {slide_number}")
        return 0.0
    print(f"✓ Slide count sufficient ({len(prs.slides)} slides)")
    score += 0.2

    # 3) Locate the title shape (0.2 points)
    slide = prs.slides[slide_number - 1]
    title_shape = slide.shapes.title
    if title_shape is None:
        # Fallback: search placeholder of TITLE type
        for shp in slide.shapes:
            if shp.is_placeholder and shp.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                title_shape = shp
                break
    if title_shape is None or not title_shape.has_text_frame:
        print("✗ Title shape with text frame not found on the slide")
        return score  # cannot continue verification
    print("✓ Title shape found")
    score += 0.2

    # 4) Analyse underline styles in all runs
    runs = [run for para in title_shape.text_frame.paragraphs for run in para.runs]
    if not runs:
        print("✗ No runs found in title text")
        return score

    double_count = single_count = none_count = 0
    for run in runs:
        underline_val = run.font.underline
        if underline_val is True or underline_val == MSO_UNDERLINE.SINGLE_LINE:
            single_count += 1
        elif underline_val == MSO_UNDERLINE.DOUBLE_LINE:
            double_count += 1
        else:
            none_count += 1

    print(f"Underline analysis -> DOUBLE: {double_count}, SINGLE: {single_count}, NONE/OTHER: {none_count}")

    # 5) Award points for at least one visible DOUBLE underline (0.4 points)
    visible_double = any(run.text.strip() and run.font.underline == MSO_UNDERLINE.DOUBLE_LINE for run in runs)
    if visible_double:
        print("✓ At least one visible run has DOUBLE underline")
        score += 0.4
    else:
        print("✗ No visible run with DOUBLE underline found")
        return score

    # 6) Extra credit if no SINGLE underline remains (0.2 points)
    if single_count == 0:
        print("✓ No SINGLE underline runs detected")
        score += 0.2
    else:
        print("✗ Some runs still have SINGLE underline")

    final_score = min(score, max_score)
    print(f"Total score: {final_score}")
    return final_score


if __name__ == "__main__":
    FILE_PATH = "/home/user/in_libreoffice_impress_im_polishing_slide_115_and_want_the_title_text_to_pop_a_bit_morespecifically__golden.pptx"
    reward = verify_double_underline_title(FILE_PATH)
    print(f"REWARD: {reward}")
