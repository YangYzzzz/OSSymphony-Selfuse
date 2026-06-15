"""
Reward Script: Academic poster (single slide, 42x36 inches) for environmental science project
Task ID: impress_stu_074
Domain: libreoffice_impress
Scoring:
  C1 - Slide dimensions 42x36 (0.15)
  C2 - Header bar with title, green fill, formatting (0.15)
  C3 - Author info line (0.10)
  C4 - Column 1: Background (5 bullets) + Research Questions (3 numbered) (0.15)
  C5 - Column 2: Methodology flowchart (5 steps) + Study Area placeholder (0.15)
  C6 - Column 3: Results (2 chart placeholders) + Conclusions (4 bullets) + Acknowledgments (0.15)
  C7 - Bottom bar with university name + QR code placeholder (0.15)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_074'

def persist_app_state(domain):
    """Save any unsaved LibreOffice state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")

def get_all_text_shapes(slide):
    """Recursively get all shapes with text frames, including grouped shapes."""
    def extract(shape):
        results = []
        if hasattr(shape, "text") and hasattr(shape, "text_frame"):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.extend(extract(sub))
        return results
    out = []
    for shape in slide.shapes:
        out.extend(extract(shape))
    return out

def find_shapes_with_text(slide, search_text, case_insensitive=True):
    """Find all shapes that contain the given text."""
    results = []
    for shape in get_all_text_shapes(slide):
        full_text = shape.text_frame.text if shape.has_text_frame else ""
        if case_insensitive:
            if search_text.lower() in full_text.lower():
                results.append(shape)
        else:
            if search_text in full_text:
                results.append(shape)
    return results

def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Inches, Emu
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    if len(prs.slides) == 0:
        print("CRITICAL: No slides in presentation")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]
    all_text_shapes = get_all_text_shapes(slide)

    # Gather all text content for later checks
    all_texts = []
    for shape in all_text_shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                txt = para.text.strip()
                if txt:
                    all_texts.append(txt)

    # Component 1: Slide dimensions — 42x36 inches (0.15 points)
    try:
        width_in = prs.slide_width / 914400
        height_in = prs.slide_height / 914400
        # Task requires 42x36 — check within 5% tolerance
        width_ok = abs(width_in - 42.0) / 42.0 < 0.05
        height_ok = abs(height_in - 36.0) / 36.0 < 0.05
        if width_ok and height_ok:
            print(f"PASS: Component 1 — Slide dimensions {width_in}x{height_in} inches (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — Expected ~42x36 inches, found {width_in}x{height_in}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Header bar with title text, green fill, 48pt white bold (0.15 points)
    try:
        c2_score = 0.0
        title_found = False
        title_green_fill = False
        title_white_bold_48pt = False

        # Find shape containing the title text
        title_shapes = find_shapes_with_text(slide, "Microplastic Contamination in Local Waterways")
        if title_shapes:
            title_found = True
            c2_score += 0.05
            ts = title_shapes[0]

            # Check green fill (#2E7D32) on the shape
            try:
                fill = ts.fill
                if fill.type == 1:  # solid
                    rgb_str = str(fill.fore_color.rgb).upper()
                    if rgb_str == "2E7D32":
                        title_green_fill = True
                        c2_score += 0.05
                    else:
                        print(f"  INFO: Header fill color is {rgb_str}, expected 2E7D32")
            except Exception:
                pass

            # Check font formatting: 48pt, white, bold
            for para in ts.text_frame.paragraphs:
                if "Microplastic Contamination" in para.text:
                    for run in para.runs:
                        if "Microplastic Contamination" in run.text:
                            font = run.font
                            size_ok = font.size is not None and abs(font.size - 609600) < 20000  # 48pt = 609600 EMU
                            bold_ok = font.bold is True
                            white_ok = False
                            try:
                                if font.color.type is not None:
                                    white_ok = str(font.color.rgb).upper() == "FFFFFF"
                            except:
                                pass
                            if size_ok and bold_ok and white_ok:
                                title_white_bold_48pt = True
                                c2_score += 0.05
                            else:
                                print(f"  INFO: Title font — size_ok={size_ok}, bold={font.bold}, white={white_ok}")
                            break
                    break

        if title_found:
            print(f"PASS: Component 2 — Header bar found (title={title_found}, green_fill={title_green_fill}, formatting={title_white_bold_48pt}) ({c2_score} pts)")
        else:
            print("FAIL: Component 2 — Title 'Microplastic Contamination in Local Waterways' not found")
        total_score += c2_score
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Author info line (0.10 points)
    try:
        author_found = False
        author_shapes = find_shapes_with_text(slide, "R. Johnson")
        for shape in author_shapes:
            for para in shape.text_frame.paragraphs:
                txt = para.text
                if "R. Johnson" in txt and "K. Patel" in txt and "ENV 350" in txt and "Dr. Williams" in txt:
                    author_found = True
                    # Check approximate font size (~22pt = 279400 EMU)
                    for run in para.runs:
                        if run.font.size is not None:
                            size_pt = run.font.size / 12700
                            print(f"  INFO: Author font size = {size_pt:.1f}pt")
                    break
            if author_found:
                break

        if author_found:
            print(f"PASS: Component 3 — Author info line found (0.10 pts)")
            total_score += 0.10
        else:
            print("FAIL: Component 3 — Author info 'R. Johnson, K. Patel | ENV 350 | Dr. Williams' not found")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Column 1 — Background section (5 bullets) + Research Questions (3 numbered) (0.15 points)
    try:
        c4_score = 0.0

        # Check for "Background" section header
        bg_headers = find_shapes_with_text(slide, "Background")
        bg_header_found = False
        for sh in bg_headers:
            # Must be a standalone header, not just text containing "Background" in a paragraph
            for para in sh.text_frame.paragraphs:
                if para.text.strip().lower() == "background":
                    bg_header_found = True
                    break
            if bg_header_found:
                break

        # Count bullet points in background section (look for shapes with bullet-like text)
        bg_bullet_count = 0
        for txt in all_texts:
            # Bullets start with bullet char or dash
            if txt.startswith("•") or txt.startswith("-") or txt.startswith("*"):
                # Background bullets contain relevant keywords
                lower_txt = txt.lower()
                if any(kw in lower_txt for kw in ["microplastic", "pollutant", "source", "waterway", "organism", "monitoring", "trophic", "fiber", "runoff"]):
                    bg_bullet_count += 1

        # Check for Research Questions section header
        rq_headers = find_shapes_with_text(slide, "Research Questions")
        rq_header_found = len(rq_headers) > 0

        # Count numbered questions
        rq_count = 0
        for txt in all_texts:
            if txt.startswith("1.") or txt.startswith("2.") or txt.startswith("3."):
                lower_txt = txt.lower()
                if any(kw in lower_txt for kw in ["microplastic", "concentration", "sampling", "correlation", "urban", "dominant"]):
                    rq_count += 1

        if bg_header_found and bg_bullet_count >= 5:
            c4_score += 0.075
        elif bg_header_found and bg_bullet_count >= 3:
            c4_score += 0.04

        if rq_header_found and rq_count >= 3:
            c4_score += 0.075
        elif rq_header_found and rq_count >= 1:
            c4_score += 0.04

        if c4_score > 0:
            print(f"PASS: Component 4 — Background ({bg_bullet_count} bullets) + Research Questions ({rq_count} questions) ({c4_score} pts)")
        else:
            print(f"FAIL: Component 4 — bg_header={bg_header_found}, bg_bullets={bg_bullet_count}, rq_header={rq_header_found}, rq_count={rq_count}")
        total_score += c4_score
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Column 2 — Methodology flowchart (5 steps) + Study Area placeholder (0.15 points)
    try:
        c5_score = 0.0

        # Check Methodology section header
        meth_headers = find_shapes_with_text(slide, "Methodology")
        meth_found = len(meth_headers) > 0

        # Count flowchart steps (Step 1 through Step 5)
        step_count = 0
        for txt in all_texts:
            lower_txt = txt.lower()
            if any(f"step {n}" in lower_txt for n in range(1, 6)):
                step_count += 1

        # Count arrow shapes (down arrows in the flowchart)
        arrow_count = 0
        for shape in slide.shapes:
            if "arrow" in shape.name.lower():
                arrow_count += 1

        # Check Study Area section
        study_headers = find_shapes_with_text(slide, "Study Area")
        study_found = len(study_headers) > 0

        # Check for map placeholder
        map_placeholder = False
        map_shapes = find_shapes_with_text(slide, "Map")
        if map_shapes:
            map_placeholder = True

        if meth_found and step_count >= 5:
            c5_score += 0.05
        elif meth_found and step_count >= 3:
            c5_score += 0.03

        if arrow_count >= 4:
            c5_score += 0.03
        elif arrow_count >= 2:
            c5_score += 0.015

        if study_found:
            c5_score += 0.035

        if map_placeholder:
            c5_score += 0.035

        if c5_score > 0:
            print(f"PASS: Component 5 — Methodology (steps={step_count}, arrows={arrow_count}), Study Area={study_found}, map={map_placeholder} ({c5_score} pts)")
        else:
            print(f"FAIL: Component 5 — meth={meth_found}, steps={step_count}, arrows={arrow_count}, study={study_found}, map={map_placeholder}")
        total_score += c5_score
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Column 3 — Results (2 chart placeholders) + Conclusions (4 bullets) + Acknowledgments (0.15 points)
    try:
        c6_score = 0.0

        # Check Results header
        results_headers = find_shapes_with_text(slide, "Results")
        results_found = False
        for sh in results_headers:
            for para in sh.text_frame.paragraphs:
                if para.text.strip().lower() == "results":
                    results_found = True
                    break
            if results_found:
                break

        # Count chart placeholders (shapes with "Chart" text)
        chart_count = 0
        for txt in all_texts:
            if "chart" in txt.lower() and ("[" in txt or "placeholder" in txt.lower()):
                chart_count += 1

        # Check Conclusions header
        conc_headers = find_shapes_with_text(slide, "Conclusions")
        conc_found = len(conc_headers) > 0

        # Count conclusion bullets
        conc_bullet_count = 0
        in_conclusions = False
        for txt in all_texts:
            if "conclusions" in txt.lower() and len(txt) < 20:
                in_conclusions = True
                continue
            if in_conclusions and (txt.startswith("•") or txt.startswith("-")):
                conc_bullet_count += 1
            # Also count bullets that follow the conclusions section by content
            if txt.startswith("•") and any(kw in txt.lower() for kw in ["concentration", "downstream", "urban", "polyester", "polyethylene", "stormwater", "intervention"]):
                if conc_bullet_count == 0:
                    conc_bullet_count += 1

        # Recount conclusions bullets more robustly: find all bullet texts after "Conclusions" header
        # by looking for shapes in column 3 area with bullet text
        conc_bullet_count_v2 = 0
        for shape in all_text_shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = para.text.strip()
                    if txt.startswith("•"):
                        # Check if this is in the right-hand third (Column 3) AND contains conclusion-like content
                        if hasattr(shape, 'left') and shape.left is not None:
                            # Column 3 starts around x=25 inches area
                            if shape.left > 20 * 914400:
                                conc_kw = ["concentration", "downstream", "urban", "polyester", "polyethylene", "stormwater", "intervention", "fiber", "particles"]
                                if any(kw in txt.lower() for kw in conc_kw):
                                    conc_bullet_count_v2 += 1

        actual_conc_bullets = max(conc_bullet_count, conc_bullet_count_v2)

        # Check Acknowledgments
        ack_shapes = find_shapes_with_text(slide, "Acknowledgments")
        ack_found = len(ack_shapes) > 0

        # Check acknowledgments has small text
        ack_small_text = False
        for sh in ack_shapes:
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is not None and run.font.size <= 177800:  # <= 14pt
                        ack_small_text = True
                        break

        if results_found and chart_count >= 2:
            c6_score += 0.05
        elif results_found and chart_count >= 1:
            c6_score += 0.03

        if conc_found and actual_conc_bullets >= 4:
            c6_score += 0.05
        elif conc_found and actual_conc_bullets >= 2:
            c6_score += 0.03

        if ack_found:
            c6_score += 0.03
        if ack_small_text:
            c6_score += 0.02

        if c6_score > 0:
            print(f"PASS: Component 6 — Results={results_found} (charts={chart_count}), Conclusions={conc_found} (bullets={actual_conc_bullets}), Acknowledgments={ack_found} ({c6_score} pts)")
        else:
            print(f"FAIL: Component 6 — results={results_found}, charts={chart_count}, conc={conc_found}, bullets={actual_conc_bullets}, ack={ack_found}")
        total_score += c6_score
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Bottom bar with university name + QR code placeholder (0.15 points)
    try:
        c7_score = 0.0

        # Find bottom bar — a shape near the bottom of the slide with green fill
        bottom_bar_found = False
        for shape in slide.shapes:
            if shape.shape_type == 1:  # AUTO_SHAPE
                # Bottom area: top position > 30 inches from top
                if shape.top is not None and shape.top > 28 * 914400:
                    try:
                        fill = shape.fill
                        if fill.type == 1:  # solid fill
                            rgb_str = str(fill.fore_color.rgb).upper()
                            if rgb_str == "2E7D32":
                                # Wide bar (spanning most of slide width)
                                if shape.width > 30 * 914400:
                                    bottom_bar_found = True
                    except:
                        pass

        # University name in bottom area
        univ_found = False
        for shape in all_text_shapes:
            if shape.has_text_frame and shape.top is not None and shape.top > 28 * 914400:
                txt = shape.text_frame.text.lower()
                if "university" in txt or "department" in txt:
                    univ_found = True
                    break

        # QR code placeholder in bottom area
        qr_found = False
        for shape in all_text_shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.lower()
                if "qr" in txt:
                    if shape.top is not None and shape.top > 28 * 914400:
                        qr_found = True
                        break

        if bottom_bar_found:
            c7_score += 0.05
        if univ_found:
            c7_score += 0.05
        if qr_found:
            c7_score += 0.05

        if c7_score > 0:
            print(f"PASS: Component 7 — Bottom bar={bottom_bar_found}, university={univ_found}, QR={qr_found} ({c7_score} pts)")
        else:
            print(f"FAIL: Component 7 — bottom_bar={bottom_bar_found}, university={univ_found}, QR={qr_found}")
        total_score += c7_score
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
