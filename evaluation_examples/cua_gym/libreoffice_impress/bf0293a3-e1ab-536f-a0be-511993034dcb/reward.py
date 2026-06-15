"""
FINAL REWARD SCRIPT - SUCCESS
Task: Shade every alternate row in Table 1 with 10% Gray (banded rows).
Generated: 2025-10-17 05:51:27
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation

# -------------------------------------------------------------
# Reward Script: Shade Every Alternate Row in Table 1 with 10% Gray
# -------------------------------------------------------------
# This script verifies the following requirements:
#   1. A table exists in the presentation (specifically, the first table found)
#   2. Every alternate row (starting with the 2nd row) is shaded Gray 10% (RGB≈E6E6E6)
#   3. Rows that are NOT supposed to be shaded remain un-filled (or at least not gray-filled)
#
# Scoring (progressive):
#   • 0.5 points → Proportion of correctly shaded rows
#   • 0.5 points → Proportion of correctly un-shaded rows
#   → Final score is the sum (capped at 1.0)
# -------------------------------------------------------------

def _get_cell_hex_rgb(cell):
    """Return the 6-digit hex RGB value of the cell's solidFill, if any.
    If no solidFill (meaning no explicit fill colour) return None."""
    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    tc_xml = cell._tc  # underlying XML of the <w:tc> element
    solid = tc_xml.find('.//a:solidFill', ns)
    if solid is None:
        return None  # no explicit fill colour → treated as un-shaded

    # solidFill can specify colour as either <a:srgbClr> or <a:schemeClr>
    srgb = solid.find('a:srgbClr', ns)
    if srgb is not None and srgb.get('val'):
        return srgb.get('val').upper()

    scheme = solid.find('a:schemeClr', ns)
    if scheme is not None:
        # "schemeClr" refers to theme colours; for banded rows this is still acceptable
        return 'SCHEME'  # treated as shaded/coloured

    return None

def _is_gray_10(hexval):
    """Return True if hex RGB looks like a light gray (~E6E6E6).
    Accept small tolerance and also treat 'SCHEME' (theme-based) as gray."""
    if hexval == 'SCHEME':
        return True
    if not hexval or len(hexval) != 6:
        return False
    try:
        r = int(hexval[0:2], 16)
        g = int(hexval[2:4], 16)
        b = int(hexval[4:6], 16)
    except ValueError:
        return False
    # Check if colour is gray (r≈g≈b) and light (~10% fill ⇒ very light gray)
    return abs(r - g) < 3 and abs(r - b) < 3 and abs(g - b) < 3 and 200 <= r <= 240

def verify_banded_rows(file_path):
    print(f"Checking presentation: {file_path}")

    # Safety: ensure file exists and is loadable
    if not os.path.exists(file_path):
        print("✗ File not found")
        print("REWARD: 0.0")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Unable to load PPTX: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Locate the first table in the deck (Table 1)
    table = None
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                break
        if table:
            break

    if table is None:
        print("✗ No table found in presentation")
        print("REWARD: 0.0")
        return 0.0

    total_rows = len(table.rows)
    print(f"✓ Table with {total_rows} rows found")

    shaded_expected = 0
    unshaded_expected = 0
    shaded_correct = 0
    unshaded_correct = 0

    for idx, row in enumerate(table.rows):
        expect_shaded = idx % 2 == 1   # 2nd, 4th, 6th… rows must be shaded
        # Inspect the first cell (representative for the row's band fill)
        cell_colour_hex = _get_cell_hex_rgb(row.cells[0])

        if expect_shaded:
            shaded_expected += 1
            if cell_colour_hex and _is_gray_10(cell_colour_hex):
                shaded_correct += 1
        else:
            unshaded_expected += 1
            # Row should remain un-filled (no solidFill) OR at least not gray
            if cell_colour_hex is None:
                unshaded_correct += 1

    # Progressive scoring
    score = 0.0
    if shaded_expected:
        score += 0.5 * (shaded_correct / shaded_expected)
    if unshaded_expected:
        score += 0.5 * (unshaded_correct / unshaded_expected)

    # Round to 3 decimals for neatness and cap at 1.0
    score = round(min(score, 1.0), 3)

    # Detailed breakdown
    print(f"Shaded rows – expected/correct: {shaded_expected}/{shaded_correct}")
    print(f"Unshaded rows – expected/correct: {unshaded_expected}/{unshaded_correct}")
    print(f"REWARD: {score}")
    return score

# -------------------------------------------------------------
# MAIN EXECUTION (update path if evaluator uses a different file)
# -------------------------------------------------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/shade_every_alternate_row_in_table_1_with_10_gray_banded_rows.pptx"
    verify_banded_rows(FILE_PATH)

