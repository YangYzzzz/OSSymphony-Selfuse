"""
FINAL REWARD SCRIPT - SUCCESS
Task: LibreOffice Impress: on slide 296, spin the image labeled “Picture 1” a full 180° so it ends up completely upside-down.
Generated: 2025-09-10 21:01:55
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation

# -----------------------------------------------------------------------------
# Reward Script: Verify that on slide 296 the image named “Picture 1” is rotated
#                a full 180° (completely upside-down).
# -----------------------------------------------------------------------------
# Scoring rubric (progressive):
#   0.0  – File/slide/shape not found or rotation incorrect
#   0.4  – Shape named “Picture 1” found on slide 296
#   1.0  – Shape found AND rotation ≈ 180° (±2° tolerance)
# -----------------------------------------------------------------------------
# IMPORTANT: This script performs *real* verification — no hard-coded success!
# -----------------------------------------------------------------------------

FILE_PATH = "/home/user/libreoffice_impress_on_slide_296_spin_the_image_labeled_picture_1_a_full_180_so_it_ends_up_completel_golden.pptx"

TOLERANCE_DEG = 2  # degrees of tolerance for rotation check
SLIDE_NUMBER   = 296  # human-readable (1-based)
SLIDE_INDEX    = SLIDE_NUMBER - 1  # 0-based index used by python-pptx
TARGET_NAME    = "picture 1"  # lower-case for comparison


def approx_equal_angle(angle, target=180, tol=TOLERANCE_DEG):
    """Return True if angle ≈ target within ±tol degrees (modulo 360)."""
    if angle is None:
        return False
    diff = abs((angle % 360) - target)
    # handle wrap-around (e.g., 358° is  -2° from 0°)
    diff = min(diff, 360 - diff)
    return diff <= tol


def verify_task(file_path: str) -> float:
    """Verify task completion and return a progressive score in [0.0, 1.0]."""
    total_score = 0.0
    print(f"Loading presentation: {file_path}")

    # ------------------------------------------------------------------
    # 1. File exists & loads (no points, prerequisite)
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File not found — task failed")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Failed to load PPTX: {e}")
        return 0.0

    # ------------------------------------------------------------------
    # 2. Slide 296 exists (no points, prerequisite)
    # ------------------------------------------------------------------
    if SLIDE_INDEX >= len(prs.slides):
        print(f"✗ Slide {SLIDE_NUMBER} not found (total slides: {len(prs.slides)})")
        return 0.0
    print("✓ Slide 296 located")
    slide = prs.slides[SLIDE_INDEX]

    # ------------------------------------------------------------------
    # 3. Locate shape named “Picture 1” — worth 0.4 points
    # ------------------------------------------------------------------
    target_shapes = []
    for shape in slide.shapes:
        name = ""
        try:
            name = shape.name or ""
        except Exception:
            pass  # some shapes may not expose .name
        if name.strip().lower() == TARGET_NAME:
            target_shapes.append(shape)

    if not target_shapes:
        print("✗ No shape named ‘Picture 1’ found on slide 296")
        return 0.0  # cannot progress without target shape

    print(f"✓ Found {len(target_shapes)} target shape(s) named ‘Picture 1’ (0.4 pts)")
    total_score += 0.4

    # ------------------------------------------------------------------
    # 4. Verify at least one of the target shapes is rotated ≈ 180° — 0.6 pts
    # ------------------------------------------------------------------
    rotation_ok = False
    for shp in target_shapes:
        angle = shp.rotation  # may be float (degrees) or None
        print(f"   • Detected rotation: {angle}°")
        if approx_equal_angle(angle):
            rotation_ok = True

    if rotation_ok:
        print("✓ Shape is rotated approximately 180° (upside-down) (+0.6 pts)")
        total_score += 0.6
    else:
        print("✗ Shape is NOT rotated 180° — no additional points")

    # ------------------------------------------------------------------
    final_score = min(total_score, 1.0)
    print(f"Total score: {final_score}")
    return final_score


if __name__ == "__main__":
    reward = verify_task(FILE_PATH)
    print(f"REWARD: {reward}")

