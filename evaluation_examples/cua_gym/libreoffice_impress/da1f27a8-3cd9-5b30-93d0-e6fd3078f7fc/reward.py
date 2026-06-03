"""
Reward Script: Set all slides to landscape orientation (10 x 7.5 inches)
Task ID: impress_fix_058
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Slide dimensions are 10 x 7.5 inches (standard landscape)
  Component 2 (0.3): All content elements fit within slide bounds
  Component 3 (0.3): All original text content is preserved
"""

import os

from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_058'

# Expected landscape dimensions in EMU
LANDSCAPE_WIDTH_EMU = 9144000   # 10 inches
LANDSCAPE_HEIGHT_EMU = 6858000  # 7.5 inches

# Tolerance for dimension check (0.5%)
DIM_TOLERANCE = 0.005

# Expected text content from initial file (all slides)
EXPECTED_TEXTS = [
    # Slide 1
    [
        "Q3 2025 Marketing Strategy",
        "Regional Performance & Growth Initiatives",
        "Prepared by: Elena Vasquez, VP Marketing",
        "August 15, 2025",
    ],
    # Slide 2
    [
        "Key Performance Metrics",
    ],
    # Slide 3
    [
        "Regional Revenue Breakdown",
        "North America",
        "$5.2M",
        "42% of total",
        "Europe",
        "$3.1M",
        "25% of total",
        "Asia Pacific",
        "$2.8M",
        "23% of total",
        "Latin America",
        "$0.8M",
        "6% of total",
        "Middle East & Africa",
        "$0.5M",
        "4% of total",
    ],
    # Slide 4
    [
        "Q3 Strategic Priorities",
        "Expand digital advertising spend by 20% in APAC markets",
        "Launch loyalty program targeting repeat customers",
        "Redesign email nurture sequences",
        "Partner with 3 new influencers",
        "A/B testing framework",
        "Reduce customer acquisition cost",
        "Develop case studies from top 5 enterprise clients",
    ],
    # Slide 5
    [
        "Implementation Timeline",
        "Week 1-2",
        "Week 3-4",
        "Week 5-6",
        "Week 7-8",
        "Week 9-10",
        "Week 11-12",
    ],
]


def is_approx_equal(val1, val2, tolerance=DIM_TOLERANCE):
    if val1 == val2:
        return True
    if val1 == 0 or val2 == 0:
        return val1 == val2
    return abs(val1 - val2) / max(abs(val1), abs(val2)) <= tolerance


def get_all_text(slide):
    """Get all text from a slide, including grouped shapes."""
    texts = []
    def extract(shape):
        if hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    t = cell.text.strip()
                    if t:
                        texts.append(t)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                extract(sub)
    for shape in slide.shapes:
        extract(shape)
    return texts


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide dimensions are 10 x 7.5 inches landscape (0.4 points)
    # This is the PRIMARY task change — portrait to landscape
    landscape_ok = False
    try:
        w = prs.slide_width
        h = prs.slide_height
        width_ok = is_approx_equal(w, LANDSCAPE_WIDTH_EMU)
        height_ok = is_approx_equal(h, LANDSCAPE_HEIGHT_EMU)
        is_landscape = w > h

        if width_ok and height_ok and is_landscape:
            landscape_ok = True
            print(f"PASS: Component 1 - Slide dimensions are 10x7.5 inches landscape "
                  f"(w={w/914400:.4f}, h={h/914400:.4f}) (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 - Expected 10x7.5 landscape, found "
                  f"w={w/914400:.4f}x{h/914400:.4f} inches, landscape={is_landscape}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Landscape AND all content elements fit within slide bounds (0.3 points)
    # Only scores if landscape orientation is correct (compound check anchored to task change)
    try:
        if not landscape_ok:
            print("FAIL: Component 2 - Skipped: slide dimensions not landscape")
        else:
            slide_w = prs.slide_width
            slide_h = prs.slide_height
            total_shapes = 0
            shapes_in_bounds = 0

            for slide in prs.slides:
                for shape in slide.shapes:
                    total_shapes += 1
                    right_edge = shape.left + shape.width
                    bottom_edge = shape.top + shape.height
                    # Allow 5% overflow tolerance for minor positioning
                    w_limit = slide_w * 1.05
                    h_limit = slide_h * 1.05
                    if right_edge <= w_limit and bottom_edge <= h_limit:
                        shapes_in_bounds += 1
                    else:
                        print(f"  INFO: Shape '{shape.name}' overflows: "
                              f"right={right_edge/914400:.2f}, bottom={bottom_edge/914400:.2f}")

            if total_shapes > 0 and shapes_in_bounds == total_shapes:
                print(f"PASS: Component 2 - Landscape + all {total_shapes} shapes in bounds (0.3 pts)")
                total_score += 0.3
            elif total_shapes > 0:
                ratio = shapes_in_bounds / total_shapes
                partial = round(0.3 * ratio, 2)
                print(f"PARTIAL: Component 2 - {shapes_in_bounds}/{total_shapes} shapes in bounds ({partial} pts)")
                total_score += partial
            else:
                print("FAIL: Component 2 - No shapes found")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Landscape AND all original text content is preserved (0.3 points)
    # Only scores if landscape orientation is correct (compound check anchored to task change)
    try:
        if not landscape_ok:
            print("FAIL: Component 3 - Skipped: slide dimensions not landscape")
        else:
            num_slides = len(prs.slides)
            if num_slides != 5:
                print(f"FAIL: Component 3 - Expected 5 slides, found {num_slides}")
            else:
                slides_with_text = 0
                for idx, slide in enumerate(prs.slides):
                    all_text = " ".join(get_all_text(slide)).lower()
                    expected = EXPECTED_TEXTS[idx]
                    found_count = 0
                    for phrase in expected:
                        if phrase.lower() in all_text:
                            found_count += 1
                        else:
                            print(f"  INFO: Slide {idx+1} missing phrase: '{phrase}'")
                    if found_count == len(expected):
                        slides_with_text += 1
                    else:
                        print(f"  INFO: Slide {idx+1}: {found_count}/{len(expected)} phrases found")

                if slides_with_text == 5:
                    print(f"PASS: Component 3 - Landscape + all text preserved across 5 slides (0.3 pts)")
                    total_score += 0.3
                else:
                    partial = round(0.3 * (slides_with_text / 5), 2)
                    print(f"PARTIAL: Component 3 - {slides_with_text}/5 slides have complete text ({partial} pts)")
                    total_score += partial
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
