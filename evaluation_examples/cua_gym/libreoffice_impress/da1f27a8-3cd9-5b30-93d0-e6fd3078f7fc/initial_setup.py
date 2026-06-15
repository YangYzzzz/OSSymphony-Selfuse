"""
Initial Setup: Create a portrait-oriented presentation with content laid out for portrait.
Task ID: impress_fix_058
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_058'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # Set portrait orientation: 7.5 x 10 inches (non-standard)
    prs.slide_width = Inches(7.5)
    prs.slide_height = Inches(10)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Title text box - centered near top, portrait width
    txBox = slide1.shapes.add_textbox(Inches(0.75), Inches(1.0), Inches(6.0), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Q3 2025 Marketing Strategy"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = "Regional Performance & Growth Initiatives"
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.name = "Arial"
    run2.font.size = Pt(20)
    run2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # Author info - positioned lower on portrait slide
    txBox2 = slide1.shapes.add_textbox(Inches(1.5), Inches(7.5), Inches(4.5), Inches(1.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p3 = tf2.paragraphs[0]
    p3.text = "Prepared by: Elena Vasquez, VP Marketing"
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.runs[0]
    run3.font.name = "Arial"
    run3.font.size = Pt(14)
    run3.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    p4 = tf2.add_paragraph()
    p4.text = "August 15, 2025"
    p4.alignment = PP_ALIGN.CENTER
    run4 = p4.runs[0]
    run4.font.name = "Arial"
    run4.font.size = Pt(12)
    run4.font.color.rgb = RGBColor(0x77, 0x77, 0x77)

    # --- Slide 2: Key Metrics Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])

    # Section title
    txTitle = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6.5), Inches(0.8))
    tf_t = txTitle.text_frame
    p_t = tf_t.paragraphs[0]
    p_t.text = "Key Performance Metrics"
    p_t.alignment = PP_ALIGN.LEFT
    r_t = p_t.runs[0]
    r_t.font.name = "Arial"
    r_t.font.size = Pt(28)
    r_t.font.bold = True
    r_t.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    # Table with metrics - portrait layout (narrow, tall)
    table_shape = slide2.shapes.add_table(
        8, 3, Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.0)
    )
    table = table_shape.table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(2.0)

    headers = ["Metric", "Q2 Actual", "Q3 Target"]
    data = [
        ["Revenue ($M)", "$12.4M", "$14.8M"],
        ["New Customers", "1,245", "1,500"],
        ["Retention Rate", "87.3%", "90.0%"],
        ["CAC", "$342", "$310"],
        ["Brand Awareness", "62%", "68%"],
        ["Email Open Rate", "24.1%", "28.0%"],
        ["Social Engagement", "3.2%", "4.0%"],
    ]

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(11)

    # --- Slide 3: Regional Breakdown ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])

    txTitle3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6.5), Inches(0.8))
    tf3 = txTitle3.text_frame
    p3t = tf3.paragraphs[0]
    p3t.text = "Regional Revenue Breakdown"
    p3t.alignment = PP_ALIGN.LEFT
    r3 = p3t.runs[0]
    r3.font.name = "Arial"
    r3.font.size = Pt(28)
    r3.font.bold = True
    r3.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    # Stacked text boxes simulating a vertical list (portrait-suited)
    regions = [
        ("North America", "$5.2M", "42% of total", RGBColor(0x2E, 0x75, 0xB6)),
        ("Europe", "$3.1M", "25% of total", RGBColor(0x4C, 0xAF, 0x50)),
        ("Asia Pacific", "$2.8M", "23% of total", RGBColor(0xFF, 0x98, 0x00)),
        ("Latin America", "$0.8M", "6% of total", RGBColor(0xE9, 0x1E, 0x63)),
        ("Middle East & Africa", "$0.5M", "4% of total", RGBColor(0x9C, 0x27, 0xB0)),
    ]

    y_pos = Inches(1.8)
    for region_name, revenue, pct, color in regions:
        # Region card - takes up portrait width
        card = slide3.shapes.add_textbox(Inches(0.75), y_pos, Inches(6.0), Inches(1.2))
        ctf = card.text_frame
        ctf.word_wrap = True

        p_name = ctf.paragraphs[0]
        p_name.text = region_name
        rn = p_name.runs[0]
        rn.font.name = "Arial"
        rn.font.size = Pt(18)
        rn.font.bold = True
        rn.font.color.rgb = color

        p_rev = ctf.add_paragraph()
        p_rev.text = f"{revenue}  |  {pct}"
        rv = p_rev.runs[0]
        rv.font.name = "Arial"
        rv.font.size = Pt(14)
        rv.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

        y_pos += Inches(1.5)

    # --- Slide 4: Strategic Priorities ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])

    txTitle4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6.5), Inches(0.8))
    tf4 = txTitle4.text_frame
    p4t = tf4.paragraphs[0]
    p4t.text = "Q3 Strategic Priorities"
    p4t.alignment = PP_ALIGN.LEFT
    r4 = p4t.runs[0]
    r4.font.name = "Arial"
    r4.font.size = Pt(28)
    r4.font.bold = True
    r4.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    priorities = [
        "Expand digital advertising spend by 20% in APAC markets",
        "Launch loyalty program targeting repeat customers (goal: 15% enrollment)",
        "Redesign email nurture sequences for higher conversion",
        "Partner with 3 new influencers in the lifestyle vertical",
        "Implement A/B testing framework for landing pages",
        "Reduce customer acquisition cost by optimizing paid search campaigns",
        "Develop case studies from top 5 enterprise clients",
    ]

    y_p = Inches(1.8)
    for i, priority in enumerate(priorities):
        tb = slide4.shapes.add_textbox(Inches(0.75), y_p, Inches(6.0), Inches(0.8))
        ptf = tb.text_frame
        ptf.word_wrap = True
        pp = ptf.paragraphs[0]
        pp.text = f"{i+1}. {priority}"
        pr = pp.runs[0]
        pr.font.name = "Arial"
        pr.font.size = Pt(14)
        pr.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        y_p += Inches(1.05)

    # --- Slide 5: Next Steps & Timeline ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])

    txTitle5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6.5), Inches(0.8))
    tf5 = txTitle5.text_frame
    p5t = tf5.paragraphs[0]
    p5t.text = "Implementation Timeline"
    p5t.alignment = PP_ALIGN.LEFT
    r5 = p5t.runs[0]
    r5.font.name = "Arial"
    r5.font.size = Pt(28)
    r5.font.bold = True
    r5.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    # Timeline items stacked vertically (portrait-suited)
    timeline = [
        ("Week 1-2", "Audit current campaigns and identify underperformers"),
        ("Week 3-4", "Launch APAC digital expansion pilot program"),
        ("Week 5-6", "Roll out loyalty program beta to select markets"),
        ("Week 7-8", "Implement A/B testing on top 10 landing pages"),
        ("Week 9-10", "Review mid-quarter metrics and adjust budgets"),
        ("Week 11-12", "Finalize Q3 reporting and prepare Q4 strategy"),
    ]

    y_tl = Inches(1.8)
    for week, desc in timeline:
        tb = slide5.shapes.add_textbox(Inches(0.75), y_tl, Inches(6.0), Inches(1.0))
        ptf = tb.text_frame
        ptf.word_wrap = True

        pw = ptf.paragraphs[0]
        pw.text = week
        rw = pw.runs[0]
        rw.font.name = "Arial"
        rw.font.size = Pt(16)
        rw.font.bold = True
        rw.font.color.rgb = RGBColor(0x2E, 0x75, 0xB6)

        pd = ptf.add_paragraph()
        pd.text = desc
        rd = pd.runs[0]
        rd.font.name = "Arial"
        rd.font.size = Pt(12)
        rd.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

        y_tl += Inches(1.3)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
