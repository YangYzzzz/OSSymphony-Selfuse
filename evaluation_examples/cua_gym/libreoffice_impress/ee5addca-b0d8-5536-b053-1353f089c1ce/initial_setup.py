"""
Initial Setup: Create master_template.pptx with default theme for brand template design
Task ID: impress_gf5_024
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import shutil
import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_024'
OUTPUT_PPTX = f'{WORKDIR}/master_template.pptx'
OUTPUT_CANONICAL = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Meridian Corp"
    slide1.placeholders[1].text = "2025 Brand Guidelines & Presentation Template"

    # --- Slide 2: Content Slide with bullet points ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "About Our Company"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Founded in 2012, Meridian Corp has grown into a global leader in sustainable technology solutions."
    p2 = body2.add_paragraph()
    p2.text = "Headquartered in San Francisco with offices in London, Tokyo, and Sydney"
    p2.level = 1
    p3 = body2.add_paragraph()
    p3.text = "Over 3,500 employees across 12 countries"
    p3.level = 1
    p4 = body2.add_paragraph()
    p4.text = "Annual revenue exceeding $2.8 billion (FY 2024)"
    p4.level = 1

    # --- Slide 3: Another content slide ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Our Core Values"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Innovation drives everything we do"
    for val in ["Sustainability in every product decision",
                "Transparency with stakeholders and partners",
                "Collaboration across teams and borders",
                "Excellence in execution and delivery"]:
        p = body3.add_paragraph()
        p.text = val
        p.level = 0

    # --- Slide 4: Two-column layout ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Key Metrics Q4 2024"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Revenue Growth: +18% YoY"
    for metric in ["Customer Retention Rate: 94.2%",
                   "Net Promoter Score: 72",
                   "Employee Satisfaction: 4.6/5.0",
                   "Carbon Offset Target: 110% achieved"]:
        p = body4.add_paragraph()
        p.text = metric
        p.level = 0

    # --- Slide 5: Blank slide for future content ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide5.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.text = "This template uses the default theme. Please customize the Master Slide before distribution."
    for run in tf.paragraphs[0].runs:
        run.font.size = Pt(14)
        run.font.italic = True
        run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    prs.save(OUTPUT_PPTX)
    print(f'Initial file created: {OUTPUT_PPTX}')

    # Also save canonical copy
    shutil.copy(OUTPUT_PPTX, OUTPUT_CANONICAL)
    print(f'Canonical copy: {OUTPUT_CANONICAL}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT_PPTX}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
