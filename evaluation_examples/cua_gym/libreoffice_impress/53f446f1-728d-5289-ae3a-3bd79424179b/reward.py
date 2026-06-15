"""
Reward Script: Apply bold + orange (#FF8C00) + underline to titles on slides 2, 3, 5
Task ID: osworld_impress_title_selective_formatting_010
Domain: libreoffice_impress
Scoring:
  - Component 1: Slide 2 title has bold + underline + orange (#FF8C00)  — 0.30 pts
  - Component 2: Slide 3 title has bold + underline + orange (#FF8C00)  — 0.30 pts
  - Component 3: Slide 5 title has bold + underline + orange (#FF8C00)  — 0.30 pts
  - Component 4: Slide 4 title remains unformatted (plain black, no bold/underline) — 0.10 pts
Total: 1.0
"""

import os

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_selective_formatting_010'
TARGET_COLOR = 'FF8C00'


def get_title_shape(slide):
    """Return the title placeholder shape for a slide, or None."""
    for shape in slide.shapes:
        if shape.has_text_frame and hasattr(shape, 'placeholder_format') \
                and shape.placeholder_format is not None \
                and shape.placeholder_format.idx == 0:
            return shape
    return None


def get_run_color(run):
    """Return hex color string (e.g. 'FF8C00') or None if no explicit RGB color."""
    try:
        if run.font.color.type is not None:
            return str(run.font.color.rgb)
    except Exception:
        pass
    return None


def check_title_formatted(slide, slide_num):
    """
    Check that ALL non-empty runs in the slide's title have:
      bold=True, underline=True, color=FF8C00
    Returns True if all conditions hold, False otherwise.
    """
    shape = get_title_shape(slide)
    if shape is None:
        print(f"FAIL: Slide {slide_num} — no title placeholder found")
        return False

    runs = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if (run.text or "").strip():
                runs.append(run)

    if not runs:
        print(f"FAIL: Slide {slide_num} — title has no non-empty runs")
        return False

    all_pass = True
    for run in runs:
        bold_ok = run.font.bold is True
        underline_ok = run.font.underline is True
        color = get_run_color(run)
        color_ok = (color == TARGET_COLOR)

        if not bold_ok:
            print(f"FAIL: Slide {slide_num} title run '{run.text}' — bold={run.font.bold} (expected True)")
            all_pass = False
        if not underline_ok:
            print(f"FAIL: Slide {slide_num} title run '{run.text}' — underline={run.font.underline} (expected True)")
            all_pass = False
        if not color_ok:
            print(f"FAIL: Slide {slide_num} title run '{run.text}' — color={color} (expected {TARGET_COLOR})")
            all_pass = False

    return all_pass


def check_title_unformatted(slide, slide_num):
    """
    Check that ALL non-empty runs in the slide's title have:
      bold is not True, underline is not True, color is black (000000) or unset
    Returns True if title is effectively plain/unformatted, False otherwise.
    """
    shape = get_title_shape(slide)
    if shape is None:
        print(f"FAIL: Slide {slide_num} — no title placeholder found (for unformatted check)")
        return False

    runs = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if (run.text or "").strip():
                runs.append(run)

    if not runs:
        # Empty title — can't definitively fail
        print(f"WARN: Slide {slide_num} — title has no non-empty runs (unformatted check)")
        return True

    all_pass = True
    for run in runs:
        # bold=False or bold=None means not bold
        if run.font.bold is True:
            print(f"FAIL: Slide {slide_num} title run '{run.text}' — bold=True (should remain unformatted)")
            all_pass = False
        # underline=False or underline=None means not underlined
        if run.font.underline is True:
            print(f"FAIL: Slide {slide_num} title run '{run.text}' — underline=True (should remain unformatted)")
            all_pass = False
        # color should remain black or unset
        color = get_run_color(run)
        if color is not None and color != '000000':
            print(f"FAIL: Slide {slide_num} title run '{run.text}' — color={color} (should be black/unset)")
            all_pass = False

    return all_pass


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have 6 slides
    num_slides = len(prs.slides)
    if num_slides < 6:
        print(f"CRITICAL: Expected 6 slides, found {num_slides}. Cannot verify.")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)

    # Component 1: Slide 2 title has bold + underline + orange (#FF8C00) (0.30 points)
    try:
        if check_title_formatted(slides[1], slide_num=2):
            print("PASS: Component 1 — Slide 2 title is bold + underline + orange FF8C00 (0.30 pts)")
            total_score += 0.30
        else:
            print("FAIL: Component 1 — Slide 2 title does not have expected bold+underline+orange formatting")
    except Exception as e:
        print(f"ERROR: Component 1 (slide 2 formatting check) — {e}")

    # Component 2: Slide 3 title has bold + underline + orange (#FF8C00) (0.30 points)
    try:
        if check_title_formatted(slides[2], slide_num=3):
            print("PASS: Component 2 — Slide 3 title is bold + underline + orange FF8C00 (0.30 pts)")
            total_score += 0.30
        else:
            print("FAIL: Component 2 — Slide 3 title does not have expected bold+underline+orange formatting")
    except Exception as e:
        print(f"ERROR: Component 2 (slide 3 formatting check) — {e}")

    # Component 3: Slide 5 title has bold + underline + orange (#FF8C00) (0.30 points)
    try:
        if check_title_formatted(slides[4], slide_num=5):
            print("PASS: Component 3 — Slide 5 title is bold + underline + orange FF8C00 (0.30 pts)")
            total_score += 0.30
        else:
            print("FAIL: Component 3 — Slide 5 title does not have expected bold+underline+orange formatting")
    except Exception as e:
        print(f"ERROR: Component 3 (slide 5 formatting check) — {e}")

    # Component 4: Slide 4 title stays plain AND at least one target slide was formatted (0.10 points)
    # This compound check verifies SELECTIVITY: the agent formatted only the right slides.
    # It requires that (a) slide 4 is unformatted, AND (b) a task change actually happened
    # (at least one of slides 2/3/5 received the formatting). This ensures initial_env scores 0.0
    # because on initial_env no slides have been formatted yet.
    try:
        any_target_formatted = False
        for idx, snum in [(1, 2), (2, 3), (4, 5)]:
            shape = get_title_shape(slides[idx])
            if shape is None:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip() and run.font.bold is True:
                        any_target_formatted = True
                        break

        slide4_plain = check_title_unformatted(slides[3], slide_num=4)

        if any_target_formatted and slide4_plain:
            print("PASS: Component 4 — Slide 4 title remains unformatted while target slides were formatted (selectivity) (0.10 pts)")
            total_score += 0.10
        elif not any_target_formatted:
            print("FAIL: Component 4 — No target slides (2/3/5) have been formatted; selectivity cannot be confirmed")
        else:
            print("FAIL: Component 4 — Slide 4 title has unexpected formatting (should remain plain)")
    except Exception as e:
        print(f"ERROR: Component 4 (slide 4 selectivity check) — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path on the VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
