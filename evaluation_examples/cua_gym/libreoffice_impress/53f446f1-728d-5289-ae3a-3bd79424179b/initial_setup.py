"""
Initial Setup: Training workshop deck with 6 slides, all titles plain/unformatted
Task ID: osworld_impress_title_selective_formatting_010
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_selective_formatting_010'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_title_plain(slide, title_text):
    """Set the title of a slide with plain black regular-weight text."""
    title_shape = slide.shapes.title
    if title_shape is None:
        return
    tf = title_shape.text_frame
    tf.clear()
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = title_text
    run.font.bold = False
    run.font.italic = False
    run.font.underline = False
    run.font.size = Pt(32)
    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # plain black


def set_content(slide, content_lines):
    """Set body text content on a slide."""
    # Find content placeholder (index 1 typically)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            for i, line in enumerate(content_lines):
                if i == 0:
                    para = tf.paragraphs[0]
                else:
                    para = tf.add_paragraph()
                para.text = line
                para.level = 0
            return


def create_initial():
    prs = Presentation()
    # Use Title+Content layout (index 1) for all slides
    layout_title_content = prs.slide_layouts[1]
    layout_title_only = prs.slide_layouts[5]

    # --- Slide 1: Title Slide (Introduction) ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    set_title_plain(slide1, "Advanced Project Management Workshop")
    # Subtitle placeholder
    for ph in slide1.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = "Professional Development Series — Q2 2025"
            break

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(layout_title_content)
    set_title_plain(slide2, "Workshop Agenda")
    set_content(slide2, [
        "Module 1: Project Planning Fundamentals",
        "Module 2: Risk Assessment and Mitigation",
        "Module 3: Stakeholder Communication",
        "Module 4: Agile Methodology in Practice",
        "Q&A and Group Discussion",
    ])

    # --- Slide 3: Module 1 – Project Planning ---
    slide3 = prs.slides.add_slide(layout_title_content)
    set_title_plain(slide3, "Project Planning Fundamentals")
    set_content(slide3, [
        "Define project scope and deliverables",
        "Establish timeline with milestones",
        "Allocate resources and assign ownership",
        "Set clear success criteria",
        "Use WBS for complex projects",
    ])

    # --- Slide 4: Module 2 – Risk Assessment (MUST remain plain) ---
    slide4 = prs.slides.add_slide(layout_title_content)
    set_title_plain(slide4, "Risk Assessment and Mitigation")
    set_content(slide4, [
        "Identify potential risks early",
        "Classify risks: High / Medium / Low",
        "Develop contingency plans",
        "Monitor risk register weekly",
        "Escalate critical risks immediately",
    ])

    # --- Slide 5: Module 3 – Stakeholder Communication ---
    slide5 = prs.slides.add_slide(layout_title_content)
    set_title_plain(slide5, "Stakeholder Communication Strategies")
    set_content(slide5, [
        "Map all stakeholders at project start",
        "Tailor communication to audience",
        "Weekly status updates via email",
        "Monthly steering committee reports",
        "Conflict resolution protocols",
    ])

    # --- Slide 6: Summary & Next Steps ---
    slide6 = prs.slides.add_slide(layout_title_content)
    set_title_plain(slide6, "Summary and Next Steps")
    set_content(slide6, [
        "Apply learnings in your current projects",
        "Complete post-workshop assessment",
        "Join the PM Community of Practice",
        "Schedule a follow-up with your manager",
        "Access resources at the internal portal",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
