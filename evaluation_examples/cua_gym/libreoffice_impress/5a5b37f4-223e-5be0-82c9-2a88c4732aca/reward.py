"""
Reward Script: SaaS Pricing Page Presentation
Task ID: impress_wf_068
Domain: libreoffice_impress
Scoring:
  C1  File exists + 6 slides                              — 0.15
  C2  Slide 1 title text (CloudSync Pro / Plans & Pricing) — 0.10
  C3  Slide 2 has 3 pricing card rectangles (Starter $29, Pro $79, Enterprise Custom) — 0.20
  C4  Slide 2 has "Most Popular" banner shape              — 0.10
  C5  Slide 2 has Appear entrance animations               — 0.10
  C6  Slide 3 has feature comparison table with check/X    — 0.10
  C7  Slide 4 has ROI calculator input rectangles          — 0.05
  C8  Slide 5 has 12 placeholder rectangles                — 0.10
  C9  Slide 6 has "Start Free Trial" CTA button            — 0.05
  C10 Blue #1976D2 used as primary color                   — 0.05
  Total: 1.0
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_068'
FILE_PATH = os.path.join(WORKDIR, 'Desktop', 'Pricing_Deck.pptx')


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: file must exist
    if not os.path.exists(file_path):
        print(f"CRITICAL: File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)

    # ── Component 1: File has exactly 6 slides (0.15 pts) ──
    try:
        num_slides = len(slides)
        if num_slides == 6:
            print(f"PASS: Component 1 — Presentation has 6 slides (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — Expected 6 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Gate: need at least 6 slides for remaining checks
    if len(slides) < 6:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # ── Component 2: Slide 1 title contains "CloudSync Pro" and "Plans & Pricing" (0.10 pts) ──
    try:
        slide1_text = ""
        for shape in slides[0].shapes:
            if shape.has_text_frame:
                slide1_text += " " + shape.text_frame.text
        slide1_lower = slide1_text.lower()
        has_cloudsync = "cloudsync pro" in slide1_lower
        has_pricing = "plans" in slide1_lower and "pricing" in slide1_lower
        if has_cloudsync and has_pricing:
            print(f"PASS: Component 2 — Slide 1 has 'CloudSync Pro' and 'Plans & Pricing' (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 — Slide 1 text: {slide1_text.strip()[:100]}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # ── Component 3: Slide 2 has 3 pricing card rectangles with correct tiers (0.20 pts) ──
    try:
        slide2 = slides[1]
        card_shapes = []
        for shape in slide2.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.has_text_frame:
                txt = shape.text_frame.text.lower()
                # Pricing cards should have tier name + price
                is_starter = "starter" in txt and "$29" in txt
                is_pro = "pro" in txt and "$79" in txt
                is_enterprise = "enterprise" in txt and "custom" in txt
                if is_starter or is_pro or is_enterprise:
                    card_shapes.append(shape)

        tiers_found = len(card_shapes)
        if tiers_found >= 3:
            print(f"PASS: Component 3 — Slide 2 has {tiers_found} pricing tier cards (0.20 pts)")
            total_score += 0.20
        elif tiers_found >= 2:
            print(f"PARTIAL: Component 3 — Found {tiers_found}/3 tier cards (0.13 pts)")
            total_score += 0.13
        elif tiers_found >= 1:
            print(f"PARTIAL: Component 3 — Found {tiers_found}/3 tier cards (0.07 pts)")
            total_score += 0.07
        else:
            print(f"FAIL: Component 3 — No pricing tier cards found on slide 2")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # ── Component 4: Slide 2 has "Most Popular" banner shape (0.10 pts) ──
    try:
        slide2 = slides[1]
        banner_found = False
        for shape in slide2.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip().lower()
                if "most popular" in txt:
                    banner_found = True
                    break
        if banner_found:
            print(f"PASS: Component 4 — 'Most Popular' banner found on slide 2 (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 4 — No 'Most Popular' banner shape found on slide 2")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # ── Component 5: Slide 2 has Appear entrance animations (0.10 pts) ──
    try:
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        has_appear_anim = False
        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide2.xml') as f:
                root = ET.parse(f).getroot()
                timing = root.find('.//p:timing', ns)
                if timing is not None:
                    # presetID="1" and presetClass="entr" = Appear entrance effect
                    ctn_elements = timing.findall('.//' + '{http://schemas.openxmlformats.org/presentationml/2006/main}cTn')
                    appear_count = 0
                    for ctn in ctn_elements:
                        preset_id = ctn.get('presetID')
                        preset_class = ctn.get('presetClass')
                        if preset_id == '1' and preset_class == 'entr':
                            appear_count += 1
                    if appear_count >= 2:
                        has_appear_anim = True

        if has_appear_anim:
            print(f"PASS: Component 5 — Appear entrance animations found on slide 2 ({appear_count} effects) (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 5 — No Appear entrance animations found on slide 2")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # ── Component 6: Slide 3 has feature comparison table with check/X marks (0.10 pts) ──
    try:
        slide3 = slides[2]
        table_found = False
        has_checks = False
        has_x_marks = False
        for shape in slide3.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_found = True
                table = shape.table
                rows = len(table.rows)
                cols = len(table.columns)
                # Look for check and X marks in cells
                for r in range(rows):
                    for c in range(cols):
                        cell_text = table.cell(r, c).text.strip()
                        if cell_text in ('\u2713', 'V', 'v', '\u2714', '\u2611'):
                            has_checks = True
                        if '\u2713' in cell_text:
                            has_checks = True
                        if cell_text in ('\u2717', '\u2718', 'X', 'x', '\u2612'):
                            has_x_marks = True
                        if '\u2717' in cell_text or '\u2718' in cell_text or '\u2719' in cell_text:
                            has_x_marks = True
                        # Also check common unicode checkmarks
                        if '\u2705' in cell_text or '\u2611' in cell_text:
                            has_checks = True
                        if '\u274c' in cell_text or '\u274e' in cell_text:
                            has_x_marks = True

        if table_found and has_checks and has_x_marks:
            print(f"PASS: Component 6 — Feature comparison table with check/X marks found (0.10 pts)")
            total_score += 0.10
        elif table_found and (has_checks or has_x_marks):
            print(f"PARTIAL: Component 6 — Table found but missing {'X marks' if not has_x_marks else 'check marks'} (0.05 pts)")
            total_score += 0.05
        elif table_found:
            print(f"PARTIAL: Component 6 — Table found but no check/X marks (0.03 pts)")
            total_score += 0.03
        else:
            print(f"FAIL: Component 6 — No feature comparison table found on slide 3")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # ── Component 7: Slide 4 has ROI calculator input field rectangles (0.05 pts) ──
    try:
        slide4 = slides[3]
        input_rects = 0
        for shape in slide4.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.has_text_frame:
                txt = shape.text_frame.text.strip().lower()
                # Input fields typically have placeholder text or are labeled rectangles
                if "enter" in txt or "input" in txt or "value" in txt or txt == "":
                    input_rects += 1
        # Also count all auto_shapes as potential input fields
        all_rects = sum(1 for s in slide4.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE)
        if input_rects >= 3 or all_rects >= 3:
            print(f"PASS: Component 7 — Slide 4 has {max(input_rects, all_rects)} input rectangles for ROI calculator (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 7 — Expected >= 3 input rectangles, found {input_rects} (auto_shapes: {all_rects})")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    # ── Component 8: Slide 5 has 12 placeholder rectangles (customer logos grid) (0.10 pts) ──
    try:
        slide5 = slides[4]
        rect_count = sum(1 for s in slide5.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE)
        if rect_count >= 12:
            print(f"PASS: Component 8 — Slide 5 has {rect_count} rectangles (>= 12 logo placeholders) (0.10 pts)")
            total_score += 0.10
        elif rect_count >= 8:
            print(f"PARTIAL: Component 8 — Slide 5 has {rect_count} rectangles (expected 12) (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 8 — Slide 5 has only {rect_count} rectangles, expected 12")
    except Exception as e:
        print(f"ERROR: Component 8 — {e}")

    # ── Component 9: Slide 6 has "Start Free Trial" CTA button (0.05 pts) ──
    try:
        slide6 = slides[5]
        cta_found = False
        for shape in slide6.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip().lower()
                if "start free trial" in txt:
                    # Check it's a shape (button-like), not just a textbox
                    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        cta_found = True
                        break
        if cta_found:
            print(f"PASS: Component 9 — 'Start Free Trial' CTA button shape found on slide 6 (0.05 pts)")
            total_score += 0.05
        else:
            # Check if text exists at all
            all_text = " ".join(s.text_frame.text for s in slide6.shapes if s.has_text_frame).lower()
            if "start free trial" in all_text:
                print(f"PARTIAL: Component 9 — 'Start Free Trial' text found but not as auto-shape button (0.02 pts)")
                total_score += 0.02
            else:
                print(f"FAIL: Component 9 — No 'Start Free Trial' CTA found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 9 — {e}")

    # ── Component 10: Blue #1976D2 used as primary color (0.05 pts) ──
    try:
        target_rgb = "1976D2"
        color_found = False
        # Check backgrounds
        for slide in slides:
            try:
                fill = slide.background.fill
                if fill.type == 1:  # SOLID
                    if str(fill.fore_color.rgb).upper() == target_rgb:
                        color_found = True
                        break
            except:
                pass

        if not color_found:
            # Check shape fills and text colors across slides
            for slide in slides:
                for shape in slide.shapes:
                    try:
                        fill = shape.fill
                        if fill.type is not None and fill.type == 1:
                            if str(fill.fore_color.rgb).upper() == target_rgb:
                                color_found = True
                                break
                    except:
                        pass
                    if color_found:
                        break
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                try:
                                    if run.font.color.type is not None:
                                        if str(run.font.color.rgb).upper() == target_rgb:
                                            color_found = True
                                            break
                                except:
                                    pass
                            if color_found:
                                break
                    if color_found:
                        break
                if color_found:
                    break

        if color_found:
            print(f"PASS: Component 10 — Blue #1976D2 found as primary color (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 10 — Blue #1976D2 not found in presentation")
    except Exception as e:
        print(f"ERROR: Component 10 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point — test against canonical artifact path
if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(FILE_PATH)
