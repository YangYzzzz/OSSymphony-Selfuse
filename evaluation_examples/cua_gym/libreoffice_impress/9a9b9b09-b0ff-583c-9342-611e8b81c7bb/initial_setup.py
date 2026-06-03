"""
Initial Setup: Art gallery slideshow with 5 slides, no auto-advance timing
Task ID: impress_tm_036
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_036'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Art gallery paintings data
    paintings = [
        {
            "title": "Starry Night Over the Rhine",
            "artist": "Vincent van Gogh",
            "year": "1888",
            "medium": "Oil on canvas",
            "description": "A luminous depiction of the night sky reflected in the waters of the Rhone river, with gas lighting from the town shimmering on the water's surface.",
            "bg_color": RGBColor(0x0A, 0x0A, 0x2E),
            "text_color": RGBColor(0xFF, 0xD7, 0x00),
        },
        {
            "title": "Water Lilies at Twilight",
            "artist": "Claude Monet",
            "year": "1906",
            "medium": "Oil on canvas",
            "description": "A serene view of Monet's beloved water garden at Giverny, capturing the interplay of light and reflection as day gives way to evening.",
            "bg_color": RGBColor(0x1B, 0x3A, 0x4B),
            "text_color": RGBColor(0xB0, 0xE0, 0xE6),
        },
        {
            "title": "The Persistence of Memory",
            "artist": "Salvador Dali",
            "year": "1931",
            "medium": "Oil on canvas",
            "description": "Soft, melting pocket watches draped over a barren landscape, exploring the fluidity and irrelevance of time in the unconscious mind.",
            "bg_color": RGBColor(0x2C, 0x1A, 0x0A),
            "text_color": RGBColor(0xF5, 0xDE, 0xB3),
        },
        {
            "title": "Girl with a Pearl Earring",
            "artist": "Johannes Vermeer",
            "year": "1665",
            "medium": "Oil on canvas",
            "description": "An enigmatic portrait of a young woman wearing an exotic turban and an oversized pearl earring, often called the 'Mona Lisa of the North'.",
            "bg_color": RGBColor(0x0F, 0x1A, 0x30),
            "text_color": RGBColor(0xFA, 0xEB, 0xD7),
        },
        {
            "title": "The Great Wave off Kanagawa",
            "artist": "Katsushika Hokusai",
            "year": "1831",
            "medium": "Woodblock print",
            "description": "A towering wave threatens boats near the coast while Mount Fuji rises serenely in the background, embodying the tension between nature and humanity.",
            "bg_color": RGBColor(0x00, 0x1F, 0x3F),
            "text_color": RGBColor(0x7F, 0xDB, 0xFF),
        },
    ]

    for i, painting in enumerate(paintings):
        # Use blank layout
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only

        # Set background color
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = painting["bg_color"]

        # Slide number label (top-right)
        txBox = slide.shapes.add_textbox(Inches(11.5), Inches(0.3), Inches(1.5), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = f"Slide {i + 1} of 5"
        run.font.size = Pt(12)
        run.font.color.rgb = painting["text_color"]
        run.font.italic = True

        # Painting title (large, centered)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = painting["title"]
        run.font.name = "Georgia"
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = painting["text_color"]

        # Artist and year
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.9), Inches(11.333), Inches(0.7))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"{painting['artist']}  |  {painting['year']}  |  {painting['medium']}"
        run.font.name = "Georgia"
        run.font.size = Pt(20)
        run.font.italic = True
        run.font.color.rgb = painting["text_color"]

        # Description
        txBox = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9.333), Inches(2.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = painting["description"]
        run.font.name = "Georgia"
        run.font.size = Pt(16)
        run.font.color.rgb = painting["text_color"]

        # Decorative separator line (text-based)
        txBox = slide.shapes.add_textbox(Inches(4), Inches(3.7), Inches(5.333), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "— — — — — — — — —"
        run.font.size = Pt(14)
        run.font.color.rgb = painting["text_color"]

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
