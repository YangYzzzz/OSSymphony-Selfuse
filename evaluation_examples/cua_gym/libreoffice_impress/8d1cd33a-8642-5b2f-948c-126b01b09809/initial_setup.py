"""
Initial Setup: Math Lecture presentation with 8 slides, titles in 24pt black.
Task ID: impress_teach_003
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_003'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_data = [
        {
            "title": "Introduction to Calculus",
            "body": "Course Overview\n\nInstructor: Dr. Elena Vasquez\nDepartment of Mathematics\nFall 2025 Semester"
        },
        {
            "title": "Limits and Continuity",
            "body": "Definition of a Limit\n\nThe limit of f(x) as x approaches a is L if for every epsilon > 0 there exists a delta > 0 such that |f(x) - L| < epsilon whenever 0 < |x - a| < delta."
        },
        {
            "title": "Derivatives and Differentiation",
            "body": "The Power Rule\n\nIf f(x) = x^n, then f'(x) = n * x^(n-1)\n\nExamples:\n  f(x) = x^3  =>  f'(x) = 3x^2\n  f(x) = x^5  =>  f'(x) = 5x^4"
        },
        {
            "title": "Applications of Derivatives",
            "body": "Finding Critical Points\n\n1. Compute f'(x)\n2. Set f'(x) = 0 and solve\n3. Classify using the second derivative test\n4. Identify local maxima and minima"
        },
        {
            "title": "Integration Fundamentals",
            "body": "The Fundamental Theorem of Calculus\n\nIf F is an antiderivative of f on [a, b], then the definite integral of f from a to b equals F(b) - F(a)."
        },
        {
            "title": "Techniques of Integration",
            "body": "Common Methods:\n\n- Substitution (u-substitution)\n- Integration by Parts\n- Partial Fraction Decomposition\n- Trigonometric Substitution"
        },
        {
            "title": "Sequences and Series",
            "body": "Convergence Tests\n\n- Ratio Test\n- Root Test\n- Comparison Test\n- Integral Test\n- Alternating Series Test"
        },
        {
            "title": "Final Review and Exam Preparation",
            "body": "Key Topics to Study:\n\n- Limits and L'Hopital's Rule\n- Differentiation Rules\n- Optimization Problems\n- Definite and Indefinite Integrals\n- Taylor and Maclaurin Series"
        },
    ]

    for i, sd in enumerate(slide_data):
        if i == 0:
            slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content

        # Set title with 24pt black font
        title_shape = slide.shapes.title
        title_shape.text = ""
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = sd["title"]
        run.font.size = Pt(24)
        run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # black
        run.font.bold = True

        # Set body content
        if i == 0:
            body_ph = slide.placeholders[1]
        else:
            body_ph = slide.placeholders[1]
        body_ph.text = sd["body"]

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
