"""
Reward Script: Change all slide title font sizes to 36pt and color to #003366
Task ID: impress_teach_003
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): All 8 slide titles have font size 36pt
  Component 2 (0.5): All 8 slide titles have font color #003366
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_003'
EXPECTED_SLIDES = 8
EXPECTED_SIZE_PT = 36.0
EXPECTED_COLOR = '003366'


def persist_app_state(domain: str):
    """Save any unsaved changes in LibreOffice before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
    except ImportError as e:
        print(f"CRITICAL: Missing python-pptx library: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has 8 slides
    num_slides = len(prs.slides)
    if num_slides != EXPECTED_SLIDES:
        print(f"PRECONDITION FAIL: Expected {EXPECTED_SLIDES} slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    # Extract title runs from all slides using slide.shapes.title (robust to shape renaming)
    title_runs = []  # list of (slide_num, run)
    for i, slide in enumerate(prs.slides):
        title_shape = slide.shapes.title
        if title_shape and title_shape.has_text_frame:
            for para in title_shape.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        title_runs.append((i + 1, run))

    if len(title_runs) == 0:
        print("FAIL: No title runs found in any slide")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Found {len(title_runs)} title run(s) across {EXPECTED_SLIDES} slides")

    # Component 1: Title font sizes are 36pt (0.5 points)
    # Progressive: each correctly-sized title contributes proportionally
    try:
        size_correct = 0
        size_total = len(title_runs)
        for slide_num, run in title_runs:
            sz = run.font.size
            if sz is not None:
                actual_pt = sz / 12700  # EMU to pt
                if abs(actual_pt - EXPECTED_SIZE_PT) < 0.5:  # tolerance of 0.5pt
                    size_correct += 1
                    print(f"  PASS: Slide {slide_num} title size = {actual_pt}pt")
                else:
                    print(f"  FAIL: Slide {slide_num} title size = {actual_pt}pt (expected {EXPECTED_SIZE_PT}pt)")
            else:
                print(f"  FAIL: Slide {slide_num} title size = None/inherited (expected {EXPECTED_SIZE_PT}pt)")

        if size_total > 0 and size_correct == size_total:
            print(f"PASS: Component 1 -- All {size_total} title runs have size {EXPECTED_SIZE_PT}pt (0.5 pts)")
            total_score += 0.5
        elif size_correct > 0:
            partial = 0.5 * (size_correct / size_total)
            print(f"PARTIAL: Component 1 -- {size_correct}/{size_total} title runs correct size ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- 0/{size_total} title runs have correct size")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Title font colors are #003366 (0.5 points)
    # Progressive: each correctly-colored title contributes proportionally
    try:
        color_correct = 0
        color_total = len(title_runs)
        for slide_num, run in title_runs:
            try:
                if run.font.color.type is not None:
                    actual_rgb = str(run.font.color.rgb).upper()
                    if actual_rgb == EXPECTED_COLOR.upper():
                        color_correct += 1
                        print(f"  PASS: Slide {slide_num} title color = #{actual_rgb}")
                    else:
                        print(f"  FAIL: Slide {slide_num} title color = #{actual_rgb} (expected #{EXPECTED_COLOR})")
                else:
                    print(f"  FAIL: Slide {slide_num} title color = None/theme (expected #{EXPECTED_COLOR})")
            except Exception as ce:
                print(f"  FAIL: Slide {slide_num} title color read error: {ce}")

        if color_total > 0 and color_correct == color_total:
            print(f"PASS: Component 2 -- All {color_total} title runs have color #{EXPECTED_COLOR} (0.5 pts)")
            total_score += 0.5
        elif color_correct > 0:
            partial = 0.5 * (color_correct / color_total)
            print(f"PARTIAL: Component 2 -- {color_correct}/{color_total} title runs correct color ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 -- 0/{color_total} title runs have correct color")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
