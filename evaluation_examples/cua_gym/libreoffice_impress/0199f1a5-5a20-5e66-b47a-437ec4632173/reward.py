"""
Reward Script: Design a poster-style slide at A1 paper size (landscape) for a scientific conference
Task ID: impress_gf2_007
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Slide dimensions are A1 landscape (84.1 x 59.4 cm)
  Component 2 (0.25): Dark navy rectangle (#0A1628) spanning left ~25% width, full height
  Component 3 (0.25): Title text box — 48pt, bold, white (#FFFFFF), within the dark panel
  Component 4 (0.25): Subtitle text box — 28pt, gray (#888888), below the title
"""

import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_007'

# Conversion constants
CM_TO_EMU = 360000  # 1 cm = 360000 EMU
PT_TO_EMU = 12700   # 1 pt = 12700 EMU


def approx_equal(val1, val2, tolerance=0.05):
    """Check if two values are approximately equal within relative tolerance."""
    if val1 == val2:
        return True
    if val1 == 0 or val2 == 0:
        return abs(val1 - val2) < 100000  # small absolute tolerance for zero
    return abs(val1 - val2) / max(abs(val1), abs(val2)) <= tolerance


def color_matches(actual_rgb, expected_hex, tolerance=5):
    """Check if an RGB color matches the expected hex, with per-channel tolerance."""
    if actual_rgb is None:
        return False
    actual_str = str(actual_rgb).upper()
    expected_hex = expected_hex.upper().lstrip('#')
    if len(actual_str) != 6 or len(expected_hex) != 6:
        return False
    for i in range(3):
        a = int(actual_str[i*2:(i+1)*2], 16)
        e = int(expected_hex[i*2:(i+1)*2], 16)
        if abs(a - e) > tolerance:
            return False
    return True


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    if len(prs.slides) == 0:
        print("CRITICAL: No slides in presentation")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]

    # =====================================================================
    # Component 1: Slide dimensions are A1 landscape (84.1 x 59.4 cm)
    # (0.25 points)
    # Initial: 25.40 x 19.05 cm  |  Golden: 84.10 x 59.40 cm
    # =====================================================================
    try:
        expected_width = 84.1 * CM_TO_EMU   # 30276000
        expected_height = 59.4 * CM_TO_EMU  # 21384000
        actual_w = prs.slide_width
        actual_h = prs.slide_height

        width_ok = approx_equal(actual_w, expected_width, tolerance=0.02)
        height_ok = approx_equal(actual_h, expected_height, tolerance=0.02)
        is_landscape = actual_w > actual_h

        if width_ok and height_ok and is_landscape:
            print(f"PASS: Component 1 — Slide is A1 landscape: {actual_w/CM_TO_EMU:.1f} x {actual_h/CM_TO_EMU:.1f} cm (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — Slide dimensions: {actual_w/CM_TO_EMU:.1f} x {actual_h/CM_TO_EMU:.1f} cm, "
                  f"expected ~84.1 x 59.4 cm landscape. width_ok={width_ok}, height_ok={height_ok}, landscape={is_landscape}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # =====================================================================
    # Component 2: Dark navy rectangle (#0A1628) spanning left ~25% width,
    # full height (0.25 points)
    # Initial: no shapes  |  Golden: rectangle with fill 0A1628
    # =====================================================================
    try:
        navy_rect_found = False
        for shape in slide.shapes:
            # Check for auto shapes (rectangles)
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or shape.shape_type == 1:
                try:
                    fill = shape.fill
                    if fill.type is not None and fill.type == 1:  # SOLID
                        fill_rgb = fill.fore_color.rgb
                        if color_matches(fill_rgb, '0A1628', tolerance=15):
                            # Check position: should be at left edge
                            left_ok = shape.left <= 0.02 * prs.slide_width
                            # Width should be ~25% of slide width (tolerance: 18-32%)
                            width_ratio = shape.width / prs.slide_width
                            width_ok = 0.18 <= width_ratio <= 0.32
                            # Height should be full slide height (tolerance: 90%+)
                            height_ratio = shape.height / prs.slide_height
                            height_ok = height_ratio >= 0.90
                            if left_ok and width_ok and height_ok:
                                navy_rect_found = True
                                print(f"PASS: Component 2 — Navy rect found: fill={fill_rgb}, "
                                      f"width_ratio={width_ratio:.2f}, height_ratio={height_ratio:.2f} (0.25 pts)")
                            else:
                                print(f"INFO: Navy-colored rect found but position/size wrong: "
                                      f"left_ok={left_ok}, width_ratio={width_ratio:.2f}, height_ratio={height_ratio:.2f}")
                except Exception:
                    pass

        if navy_rect_found:
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 — No dark navy (#0A1628) rectangle found spanning left ~25% of slide")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # =====================================================================
    # Component 3: Title text box — 48pt, bold, white (#FFFFFF), within
    # the dark panel area (0.25 points)
    # Initial: no shapes  |  Golden: textbox with 48pt bold white text
    # =====================================================================
    try:
        title_found = False
        panel_right_edge = prs.slide_width * 0.32  # generous panel boundary

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    runs = [r for r in para.runs if (r.text or "").strip()]
                    if not runs:
                        continue
                    for run in runs:
                        font_size = run.font.size
                        is_bold = run.font.bold is True
                        try:
                            color_rgb = run.font.color.rgb if run.font.color.type is not None else None
                        except:
                            color_rgb = None

                        # Check: ~48pt, bold, white, positioned within the left panel
                        if font_size is not None:
                            size_pt = font_size / PT_TO_EMU
                            is_48pt = approx_equal(size_pt, 48, tolerance=0.05)
                            is_white = color_matches(color_rgb, 'FFFFFF', tolerance=10)
                            in_panel = shape.left < panel_right_edge

                            if is_48pt and is_bold and is_white and in_panel:
                                title_found = True
                                print(f"PASS: Component 3 — Title found: \"{run.text[:50]}\", "
                                      f"size={size_pt:.0f}pt, bold={is_bold}, color={color_rgb} (0.25 pts)")
                                break
                    if title_found:
                        break
            if title_found:
                break

        if title_found:
            total_score += 0.25
        else:
            print("FAIL: Component 3 — No title textbox found with 48pt bold white text in the left panel")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # =====================================================================
    # Component 4: Subtitle text box — 28pt, gray (#888888), below the
    # title (0.25 points)
    # Initial: no shapes  |  Golden: textbox with 28pt gray text
    # =====================================================================
    try:
        subtitle_found = False

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    runs = [r for r in para.runs if (r.text or "").strip()]
                    if not runs:
                        continue
                    for run in runs:
                        font_size = run.font.size
                        try:
                            color_rgb = run.font.color.rgb if run.font.color.type is not None else None
                        except:
                            color_rgb = None

                        if font_size is not None:
                            size_pt = font_size / PT_TO_EMU
                            is_28pt = approx_equal(size_pt, 28, tolerance=0.05)
                            is_gray = color_matches(color_rgb, '888888', tolerance=15)

                            if is_28pt and is_gray:
                                subtitle_found = True
                                print(f"PASS: Component 4 — Subtitle found: \"{run.text[:50]}\", "
                                      f"size={size_pt:.0f}pt, color={color_rgb} (0.25 pts)")
                                break
                    if subtitle_found:
                        break
            if subtitle_found:
                break

        if subtitle_found:
            total_score += 0.25
        else:
            print("FAIL: Component 4 — No subtitle textbox found with 28pt gray (#888888) text")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
