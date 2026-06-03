"""
Initial Setup: Remove all bullets from the list on slide 4
Task ID: impstruct_028
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impstruct_028'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_bullet_paragraph(tf, text, bullet_char='\u2022', level=0, font_size=Pt(18), is_first=False):
    """Add a bulleted paragraph to a text frame."""
    if is_first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()

    p.level = level
    p.space_before = Pt(6)
    p.space_after = Pt(6)

    # Set bullet character via XML
    pPr = p._p.get_or_add_pPr()
    pPr.set('lvl', str(level))
    pPr.set('indent', str(-228600))  # negative indent for bullet hang
    pPr.set('marL', str(457200))     # left margin

    bu_char = pPr.makeelement(qn('a:buChar'), {'char': bullet_char})
    pPr.append(bu_char)

    bu_sz = pPr.makeelement(qn('a:buSzPct'), {'val': '100000'})
    pPr.append(bu_sz)

    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.name = 'Calibri'
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    return p


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q2 2025 Strategy Meeting"
    slide1.placeholders[1].text = "Product & Engineering Review\nApril 14, 2025"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide2.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Agenda"
    p.runs[0].font.size = Pt(28)
    p.runs[0].font.bold = True
    p.runs[0].font.name = 'Calibri'
    p.runs[0].font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    agenda_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8), Inches(5))
    tf2 = agenda_box.text_frame
    tf2.word_wrap = True
    agenda_items = [
        "Product roadmap updates",
        "Engineering velocity review",
        "Discussion points from stakeholders",
        "Action items and next steps",
    ]
    for i, item in enumerate(agenda_items):
        add_bullet_paragraph(tf2, item, font_size=Pt(20), is_first=(i == 0))

    # --- Slide 3: Product Updates ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(0.8))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Product Updates"
    p3.runs[0].font.size = Pt(28)
    p3.runs[0].font.bold = True
    p3.runs[0].font.name = 'Calibri'
    p3.runs[0].font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    content3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8), Inches(5))
    tf3c = content3.text_frame
    tf3c.word_wrap = True
    updates = [
        "Mobile app v3.2 shipped with 98.7% crash-free sessions",
        "Dashboard redesign A/B test showing +14% engagement",
        "API latency reduced from 320ms to 185ms after cache migration",
        "Customer onboarding flow rework targeting May 1 release",
    ]
    for i, item in enumerate(updates):
        add_bullet_paragraph(tf3c, item, font_size=Pt(18), is_first=(i == 0))

    # --- Slide 4: Discussion Points (THE TARGET SLIDE) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(0.8))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Discussion Points"
    p4.runs[0].font.size = Pt(28)
    p4.runs[0].font.bold = True
    p4.runs[0].font.name = 'Calibri'
    p4.runs[0].font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    discussion_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8), Inches(5))
    tf4c = discussion_box.text_frame
    tf4c.word_wrap = True
    discussion_items = [
        "Should we prioritize the enterprise SSO integration over the self-serve analytics feature for Q3?",
        "The current sprint velocity has dropped 18% since March \u2014 we need to identify root causes and staffing gaps.",
        "Partner API rate limits are affecting three key customers; consider upgrading our tier or implementing request batching.",
        "Legal has flagged GDPR compliance gaps in the data export module \u2014 remediation timeline needs to be agreed upon.",
    ]
    for i, item in enumerate(discussion_items):
        add_bullet_paragraph(tf4c, item, font_size=Pt(18), is_first=(i == 0))

    # --- Slide 5: Action Items ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox5 = slide5.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(0.8))
    tf5 = txBox5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Action Items"
    p5.runs[0].font.size = Pt(28)
    p5.runs[0].font.bold = True
    p5.runs[0].font.name = 'Calibri'
    p5.runs[0].font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    actions_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8), Inches(5))
    tf5c = actions_box.text_frame
    tf5c.word_wrap = True
    actions = [
        "Elena to finalize Q3 feature priority matrix by April 21",
        "DevOps team to run capacity planning exercise this sprint",
        "Raj to schedule partner API upgrade call with vendor",
        "Sarah to coordinate GDPR audit with legal by end of month",
    ]
    for i, item in enumerate(actions):
        add_bullet_paragraph(tf5c, item, font_size=Pt(18), is_first=(i == 0))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
