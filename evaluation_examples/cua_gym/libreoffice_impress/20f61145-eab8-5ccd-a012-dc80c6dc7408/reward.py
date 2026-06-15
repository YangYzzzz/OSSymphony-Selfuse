"""
Reward Script: Remove all bullets from the list on slide 4, converting to plain paragraphs.
Task ID: impstruct_028
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): Bullets removed from all 4 paragraphs AND text content preserved
  Component 2 (0.5): Indentation reset to zero (marL=0, indent=0) for all 4 paragraphs
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impstruct_028'

# Expected paragraph texts (must be preserved -- used as gate, not scoring)
EXPECTED_TEXTS = [
    'Should we prioritize the enterprise SSO integration over the self-serve analytics feature for Q3?',
    'The current sprint velocity has dropped 18% since March \u2014 we need to identify root causes and staffing gaps.',
    'Partner API rate limits are affecting three key customers; consider upgrading our tier or implementing request batching.',
    'Legal has flagged GDPR compliance gaps in the data export module \u2014 remediation timeline needs to be agreed upon.',
]


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.oxml.ns import qn
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file has at least 4 slides
    if len(prs.slides) < 4:
        print(f"CRITICAL: Expected at least 4 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)

    # Find the text box with the bulleted list (the content text box with 4 paragraphs)
    target_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame:
            paras = [p for p in shape.text_frame.paragraphs if p.text.strip()]
            if len(paras) == 4:
                target_shape = shape
                break

    if target_shape is None:
        print("CRITICAL: Could not find the content text box with 4 paragraphs on slide 4")
        print("REWARD: 0.0")
        return 0.0

    content_paras = [p for p in target_shape.text_frame.paragraphs if p.text.strip()]
    print(f"Found target shape with {len(content_paras)} content paragraphs")

    # Precondition gate: text content must be preserved (not a scoring component)
    texts_mismatch_count = 0
    for i, para in enumerate(content_paras):
        actual_text = para.text.strip()
        if i < len(EXPECTED_TEXTS) and actual_text != EXPECTED_TEXTS[i]:
            texts_mismatch_count += 1
            print(f"GATE FAIL: Para {i} text changed. Expected: {EXPECTED_TEXTS[i][:50]}..., Got: {actual_text[:50]}...")
    if texts_mismatch_count > 0:
        print("GATE: Text content was corrupted, cannot award points")
        print("REWARD: 0.0")
        return 0.0
    print("GATE: Text content preserved -- proceeding with verification")

    # Component 1: Bullets removed from all 4 paragraphs (0.5 points)
    # In initial: buChar='bullet char' is present. In golden: buNone=True and no buChar.
    try:
        bullets_removed_count = 0
        for i, para in enumerate(content_paras):
            pPr = para._p.find(qn('a:pPr'))
            if pPr is not None:
                buChar = pPr.find(qn('a:buChar'))
                buAutoNum = pPr.find(qn('a:buAutoNum'))

                # Bullet is removed if no buChar and no buAutoNum
                has_bullet = (buChar is not None) or (buAutoNum is not None)
                if not has_bullet:
                    bullets_removed_count += 1
                    print(f"  Para {i}: bullet removed")
                else:
                    print(f"  FAIL Para {i}: bullet still present (buChar={buChar is not None}, buAutoNum={buAutoNum is not None})")
            else:
                # No pPr means no bullet formatting
                bullets_removed_count += 1
                print(f"  Para {i}: no pPr element, no bullets")

        if bullets_removed_count == 4:
            print(f"PASS: Component 1 -- All 4 bullets removed (0.5 pts)")
            total_score += 0.5
        elif bullets_removed_count > 0:
            partial = 0.5 * (bullets_removed_count / 4)
            print(f"PARTIAL: Component 1 -- {bullets_removed_count}/4 bullets removed ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- No bullets removed")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Indentation reset to zero for all 4 paragraphs (0.5 points)
    # In initial: marL=457200, indent=-228600. In golden: marL=0, indent=0.
    try:
        indent_reset_count = 0
        for i, para in enumerate(content_paras):
            pPr = para._p.find(qn('a:pPr'))
            if pPr is not None:
                marL = pPr.get('marL')
                indent = pPr.get('indent')
                marL_val = int(marL) if marL is not None else 0
                indent_val = int(indent) if indent is not None else 0
                if marL_val == 0 and indent_val == 0:
                    indent_reset_count += 1
                    print(f"  Para {i}: indentation reset (marL={marL_val}, indent={indent_val})")
                else:
                    print(f"  FAIL Para {i}: indentation not reset (marL={marL_val}, indent={indent_val})")
            else:
                # No pPr means no indentation
                indent_reset_count += 1
                print(f"  Para {i}: no pPr, no indentation")

        if indent_reset_count == 4:
            print(f"PASS: Component 2 -- All 4 paragraphs have zero indentation (0.5 pts)")
            total_score += 0.5
        elif indent_reset_count > 0:
            partial = 0.5 * (indent_reset_count / 4)
            print(f"PARTIAL: Component 2 -- {indent_reset_count}/4 paragraphs reset ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 -- No indentation reset")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist app state before verification (LibreOffice may have unsaved changes)
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
