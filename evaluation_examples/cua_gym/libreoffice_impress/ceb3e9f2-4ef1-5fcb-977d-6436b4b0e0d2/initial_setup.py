"""
Initial Setup: Construction Update Presentation (7 slides)
Task ID: impress_tm_021
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_021'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=None, alignment=None):
    """Helper to add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14):
    """Add a bulleted list to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(font_size)
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===================== Slide 1: Title Slide =====================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    add_text_box(slide1, Inches(1.5), Inches(1.5), Inches(10), Inches(1.5),
                 "Riverdale Commercial Complex", font_size=40, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, Inches(1.5), Inches(3.2), Inches(10), Inches(1),
                 "Construction Progress Update — Q1 2026", font_size=24,
                 color=RGBColor(0xCC, 0xDD, 0xEE), alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, Inches(1.5), Inches(5.0), Inches(10), Inches(0.8),
                 "Prepared by: Elena Vasquez, Senior Project Manager\nApril 2, 2026",
                 font_size=16, color=RGBColor(0x99, 0xBB, 0xDD), alignment=PP_ALIGN.CENTER)

    # ===================== Slide 2: Project Overview =====================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide2, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "Project Overview", font_size=32, bold=True,
                 color=RGBColor(0x1B, 0x3A, 0x5C))
    add_bullet_list(slide2, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5), [
        "Total project budget: $47.8 million",
        "Current spend: $18.3 million (38.3% of budget)",
        "Timeline: 24 months (started August 2025)",
        "Overall completion: 32% — on track",
        "Zero lost-time incidents in past 90 days",
        "Workforce: 187 active personnel on site",
        "Next milestone: structural steel for Tower B (April 15)",
    ], font_size=16)

    # Add a small table for key dates
    table_shape = slide2.shapes.add_table(5, 3, Inches(7), Inches(1.8), Inches(5.5), Inches(3))
    table = table_shape.table
    headers = ["Milestone", "Target Date", "Status"]
    rows_data = [
        ["Foundation Complete", "Dec 2025", "Done"],
        ["Demolition Phase", "Mar 2026", "In Progress"],
        ["Steel Erection", "Jun 2026", "Upcoming"],
        ["Envelope Closure", "Nov 2026", "Planned"],
    ]
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
        for run in table.cell(0, c).text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(12)
    for r, row in enumerate(rows_data, 1):
        for c, val in enumerate(row):
            table.cell(r, c).text = val
            for run in table.cell(r, c).text_frame.paragraphs[0].runs:
                run.font.size = Pt(11)

    # ===================== Slide 3: Site Preparation =====================
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide3, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "Site Preparation & Earthwork", font_size=32, bold=True,
                 color=RGBColor(0x1B, 0x3A, 0x5C))
    add_bullet_list(slide3, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5), [
        "Grading and leveling: 100% complete",
        "Soil stabilization: Lime treatment applied to sectors C, D",
        "Erosion control barriers installed along south perimeter",
        "Underground utility relocation: gas and telecom lines rerouted",
        "Stormwater detention basin excavated (capacity: 42,000 gallons)",
        "Environmental monitoring wells placed at 6 locations",
    ], font_size=16)
    add_bullet_list(slide3, Inches(7), Inches(1.5), Inches(5.5), Inches(5), [
        "Geotechnical report confirms bearing capacity at 4,500 psf",
        "Dewatering system operational since January 2026",
        "Compaction testing: all sectors pass 95% Proctor density",
        "Temporary access road widened to accommodate crane mobilization",
        "Topsoil stockpiled (2,800 cubic yards) for final landscaping",
    ], font_size=16)

    # ===================== Slide 4: Demolition Phase =====================
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide4, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "Demolition Phase", font_size=32, bold=True,
                 color=RGBColor(0x1B, 0x3A, 0x5C))
    add_bullet_list(slide4, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5), [
        "Former warehouse structure (Building A): fully demolished",
        "Concrete crushing on-site: 3,200 tons recycled as aggregate",
        "Asbestos abatement completed in sectors A-1 and A-3",
        "Underground storage tanks removed (2 x 10,000 gal diesel)",
        "Soil remediation in progress for petroleum-impacted area",
        "Lead paint containment: all debris disposed at certified facility",
    ], font_size=16)
    add_bullet_list(slide4, Inches(7), Inches(1.5), Inches(5.5), Inches(5), [
        "Building B partial demolition: 60% complete",
        "Selective salvage of structural steel: 48 tons recovered",
        "Noise mitigation: sound barriers installed along Oak Street",
        "Dust suppression via water trucks (3 units, 8-hour rotation)",
        "Demolition permit renewal filed — expected approval April 10",
    ], font_size=16)

    # ===================== Slide 5: Foundation Work =====================
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide5, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "Foundation & Structural Work", font_size=32, bold=True,
                 color=RGBColor(0x1B, 0x3A, 0x5C))
    add_bullet_list(slide5, Inches(0.8), Inches(1.5), Inches(11), Inches(5), [
        "Tower A mat foundation: poured and cured (28-day strength verified at 5,200 psi)",
        "Tower B pile driving: 120 of 180 H-piles driven to refusal",
        "Grade beams for parking structure: forming in progress",
        "Rebar placement for retaining wall: Level 1 completed, Level 2 started",
        "Waterproofing membrane applied to Tower A foundation walls",
        "Concrete supplier: Apex Ready-Mix — 3 trucks daily average",
        "Post-tension cables for Tower A Level 1 slab: installation begins April 7",
    ], font_size=16)

    # ===================== Slide 6: Safety Report =====================
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide6, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "Safety & Compliance Report", font_size=32, bold=True,
                 color=RGBColor(0x1B, 0x3A, 0x5C))

    # Safety metrics table
    table_shape2 = slide6.shapes.add_table(7, 2, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.5))
    tbl2 = table_shape2.table
    safety_data = [
        ["Metric", "Value"],
        ["Days without LTI", "94"],
        ["Near-miss reports (Q1)", "12"],
        ["Toolbox talks conducted", "63"],
        ["Safety inspections", "28"],
        ["PPE compliance rate", "99.2%"],
        ["OSHA recordable rate", "0.8"],
    ]
    for r, row in enumerate(safety_data):
        for c, val in enumerate(row):
            tbl2.cell(r, c).text = val
            for run in tbl2.cell(r, c).text_frame.paragraphs[0].runs:
                run.font.size = Pt(13)
                if r == 0:
                    run.font.bold = True

    add_bullet_list(slide6, Inches(7), Inches(1.5), Inches(5.5), Inches(4.5), [
        "Weekly crane inspection logs up to date",
        "Fall protection audit passed — all harnesses within certification",
        "Emergency evacuation drill conducted March 18",
        "New traffic management plan approved for Oak Street detour",
    ], font_size=16)

    # ===================== Slide 7: Next Steps =====================
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    fill7 = slide7.background.fill
    fill7.solid()
    fill7.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    add_text_box(slide7, Inches(1.5), Inches(1.0), Inches(10), Inches(1),
                 "Next Steps & Upcoming Milestones", font_size=36, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_bullet_list(slide7, Inches(1.5), Inches(2.5), Inches(10), Inches(4), [
        "April 7: Begin post-tension cable installation (Tower A, Level 1)",
        "April 10: Demolition permit renewal expected",
        "April 15: Structural steel delivery for Tower B",
        "April 22: Third-party concrete testing — Tower A slab",
        "May 1: Crane #2 mobilization for Tower B steel erection",
        "May 15: Owner's monthly progress review meeting",
        "June 1: Target start for Tower B above-grade construction",
    ], font_size=18)

    # Fix bullet list text color on dark slide 7
    for shape in slide7.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        current_rgb = run.font.color.rgb
                        if current_rgb is None or str(current_rgb) == '000000':
                            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    except AttributeError:
                        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
