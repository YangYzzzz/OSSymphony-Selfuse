"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 216 is still using the default theme. In LibreOffice Impress, how do I change just that slide so the background is a solid #000000 and the title text is #FFFFFF?
Generated: 2025-09-10 19:59:57
Status: success
Model: azure-o3
Total Steps: 1
"""

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
import os

def verify_slide_background_and_title_color(file_path: str) -> float:
    """
    Verifies that slide 216 in the given presentation has:
      1. A solid black (#000000) background.
      2. Title text whose runs are all white (#FFFFFF).

    Returns a progressive score between 0.0 and 1.0.
    """
    print(f"Starting verification for: {file_path}")

    # Initial sanity checks (no score for these – natural conditions)
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded successfully with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Ensure slide 216 exists (index 215)
    if len(prs.slides) < 216:
        print("✗ Presentation has fewer than 216 slides")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[215]  # zero-based index

    total_score = 0.0  # progressive scoring
    max_score = 1.0

    # ---------------------------------------------
    # Requirement 1: Background is solid #000000
    # ---------------------------------------------
    bg_ok = False
    try:
        bg_fill = slide.background.fill
        if bg_fill.type == MSO_FILL.SOLID:
            # rgb returns an RGBColor object or None
            rgb = bg_fill.fore_color.rgb
            print(f"Background fill type: SOLID, RGB: {rgb}")
            if rgb and str(rgb).lower() == '000000':
                bg_ok = True
    except Exception as e:
        print(f"Error checking background: {e}")

    if bg_ok:
        print("✓ Slide 216 background is solid #000000 (0.5 points)")
        total_score += 0.5
    else:
        print("✗ Background is not solid #000000")

    # ---------------------------------------------
    # Requirement 2: Title text is #FFFFFF
    # ---------------------------------------------
    title_ok = False
    title_shape = None

    # Locate the title placeholder
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
            title_shape = shape
            break

    if title_shape is None:
        print("✗ No title placeholder found on slide 216")
    else:
        all_white = True
        any_text = False
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text.strip():
                    any_text = True
                    run_color = run.font.color.rgb
                    print(f"  Found run '{run.text}' with color: {run_color}")
                    if run_color is None or str(run_color).lower() != 'ffffff':
                        all_white = False
        if any_text and all_white:
            title_ok = True

    if title_ok:
        print("✓ All title text on slide 216 is #FFFFFF (0.5 points)")
        total_score += 0.5
    else:
        if title_shape is not None:
            print("✗ Title text is not uniformly #FFFFFF or missing text runs")

    # ---------------------------------------------
    # Final score
    # ---------------------------------------------
    final_score = min(total_score, max_score)
    print(f"REWARD: {final_score}")
    return final_score


# -----------------------------------------------------------------------------
# Main execution for the reward script
# -----------------------------------------------------------------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/slide_216_is_still_using_the_default_theme_in_libreoffice_impress_how_do_i_change_just_that_slide_so_golden.pptx"
    verify_slide_background_and_title_color(FILE_PATH)

