"""
Reward Script: Fix transparent image on slide 4 by adding white rectangle behind it
Task ID: impress_fix_057
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): White-filled rectangle exists on slide 4
  Component 2 (0.3): Rectangle is behind the image in z-order
  Component 3 (0.3): Rectangle covers the image area (position/size match)
"""

import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_057'


def persist_app_state(domain):
    """Save any unsaved LibreOffice changes before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0

    Task: Add a white filled rectangle behind the PNG image on slide 4
    to fix transparency issues with the dark background.
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # slide 4 (0-indexed)

    # Find all shapes on slide 4, categorized
    picture_shapes = []
    white_rect_shapes = []
    all_shapes = list(slide.shapes)

    for i, shape in enumerate(all_shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture_shapes.append((i, shape))
        elif shape.shape_type == 1:  # AUTO_SHAPE (rectangle)
            try:
                fill = shape.fill
                if fill.type is not None and fill.type == 1:  # SOLID fill
                    rgb = fill.fore_color.rgb
                    if str(rgb).upper() == 'FFFFFF':
                        white_rect_shapes.append((i, shape))
            except Exception:
                pass

    # Component 1: White-filled rectangle exists on slide 4 (0.4 points)
    # This check FAILS on initial (no white rect) and PASSES on golden (white rect added)
    try:
        if len(white_rect_shapes) > 0:
            print(f"PASS: Component 1 -- Found {len(white_rect_shapes)} white-filled rectangle(s) on slide 4 (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 -- No white-filled rectangle found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: The white rectangle is behind the image in z-order (0.3 points)
    # Z-order in python-pptx = index in slide.shapes; lower index = further back
    # This FAILS on initial (no white rect at all) and PASSES on golden (rect before image)
    try:
        if len(white_rect_shapes) > 0 and len(picture_shapes) > 0:
            # Check if any white rectangle has a lower z-index than any picture
            behind_pairs = [
                (rect_idx, rect_shape, pic_idx, pic_shape)
                for rect_idx, rect_shape in white_rect_shapes
                for pic_idx, pic_shape in picture_shapes
                if rect_idx < pic_idx
            ]

            if len(behind_pairs) > 0:
                ri, rs, pi, ps = behind_pairs[0]
                print(f"  Detail: Rect '{rs.name}' z={ri} behind image '{ps.name}' z={pi}")
                total_score += 0.3
                print(f"PASS: Component 2 -- White rectangle is behind the image in z-order (0.3 pts)")
            else:
                print(f"FAIL: Component 2 -- White rectangle exists but is NOT behind the image in z-order")
        else:
            if len(white_rect_shapes) == 0:
                print(f"FAIL: Component 2 -- No white rectangle found (prerequisite failed)")
            else:
                print(f"FAIL: Component 2 -- No image found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: The white rectangle covers the image area (0.3 points)
    # The rectangle should be at least as large as and overlapping the image position
    # This FAILS on initial (no white rect) and PASSES on golden (rect matches image area)
    try:
        if len(white_rect_shapes) > 0 and len(picture_shapes) > 0:
            # Find the best matching white rect for the picture
            tolerance = 0.1  # 10% tolerance for position/size comparison
            found_coverage = 0  # 0 = not found, 1 = found

            for _, rect_shape in white_rect_shapes:
                for _, pic_shape in picture_shapes:
                    rect_left = rect_shape.left
                    rect_top = rect_shape.top
                    rect_right = rect_left + rect_shape.width
                    rect_bottom = rect_top + rect_shape.height

                    pic_left = pic_shape.left
                    pic_top = pic_shape.top
                    pic_right = pic_left + pic_shape.width
                    pic_bottom = pic_top + pic_shape.height

                    # Allow rectangle to be slightly larger or same size
                    width_tol = pic_shape.width * tolerance
                    height_tol = pic_shape.height * tolerance

                    left_ok = rect_left <= pic_left + width_tol
                    top_ok = rect_top <= pic_top + height_tol
                    right_ok = rect_right >= pic_right - width_tol
                    bottom_ok = rect_bottom >= pic_bottom - height_tol

                    print(f"  Detail: Rect pos=({rect_left},{rect_top}) size=({rect_shape.width},{rect_shape.height})")
                    print(f"  Detail: Image pos=({pic_left},{pic_top}) size=({pic_shape.width},{pic_shape.height})")
                    print(f"  Detail: Coverage: left={left_ok}, top={top_ok}, right={right_ok}, bottom={bottom_ok}")

                    if left_ok and top_ok and right_ok and bottom_ok:
                        found_coverage = 1
                        break
                if found_coverage > 0:
                    break

            if found_coverage > 0:
                print(f"PASS: Component 3 -- White rectangle covers the image area (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 3 -- White rectangle does not sufficiently cover the image area")
        else:
            print(f"FAIL: Component 3 -- Missing white rectangle or image on slide 4")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
