"""
Initial Setup: Add white rectangle behind transparent PNG on slide 4
Task ID: impress_fix_057
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from io import BytesIO

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_057'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_transparent_png():
    """Create a PNG image with transparency (a product logo with transparent background)."""
    img = Image.new('RGBA', (400, 300), (0, 0, 0, 0))  # fully transparent
    from PIL import ImageDraw, ImageFont
    draw = ImageDraw.Draw(img)

    # Draw a rounded rectangle as logo shape (white fill with some detail)
    # Main logo block
    draw.rounded_rectangle([40, 30, 360, 200], radius=20, fill=(255, 255, 255, 255))
    # Inner accent bar
    draw.rectangle([60, 50, 340, 80], fill=(52, 120, 198, 255))
    # Text area
    draw.text((80, 100), "TechNova", fill=(33, 33, 33, 255))
    draw.text((80, 140), "PRODUCT SPEC", fill=(100, 100, 100, 255))

    # Bottom detail line (partially transparent)
    draw.rectangle([40, 220, 360, 240], fill=(52, 120, 198, 180))
    draw.text((50, 222), "v3.2 | REV-2025", fill=(255, 255, 255, 200))

    # Some dots that rely on transparency
    for x in range(60, 350, 30):
        draw.ellipse([x, 260, x + 15, 275], fill=(80, 80, 80, 120))

    buf = BytesIO()
    img.save(buf, format='PNG')
    buf.seek(0)
    return buf


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    bg1 = slide1.background.fill
    bg1.solid()
    bg1.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)

    tx1 = slide1.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(2))
    tf1 = tx1.text_frame
    p = tf1.paragraphs[0]
    p.text = "TechNova Pro X200"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    p2 = tf1.add_paragraph()
    p2.text = "Product Specification Sheet"
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.runs[0]
    r2.font.size = Pt(24)
    r2.font.color.rgb = RGBColor(0x7E, 0xB8, 0xDA)

    p3 = tf1.add_paragraph()
    p3.text = "Revision 3.2 | March 2025"
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.runs[0]
    r3.font.size = Pt(16)
    r3.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)

    # --- Slide 2: Key Features ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    bg2 = slide2.background.fill
    bg2.solid()
    bg2.fore_color.rgb = RGBColor(0x12, 0x23, 0x35)

    tx_title2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(1))
    tf_t2 = tx_title2.text_frame
    p_t2 = tf_t2.paragraphs[0]
    p_t2.text = "Key Features"
    r_t2 = p_t2.runs[0]
    r_t2.font.size = Pt(32)
    r_t2.font.bold = True
    r_t2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    features = [
        ("12-Core ARM Cortex-A78 Processor", "Delivers up to 3.8 GHz burst clock for intensive workloads"),
        ("64 GB LPDDR5X Memory", "6400 MT/s bandwidth with ECC support for enterprise reliability"),
        ("2 TB NVMe Gen5 Storage", "Sequential read speeds up to 14,000 MB/s"),
        ("Wi-Fi 7 + Bluetooth 5.4", "Tri-band connectivity with 320 MHz channel support"),
        ("100W USB-C PD Charging", "Fast charge from 0 to 80% in 35 minutes"),
    ]
    y_pos = 1.6
    for title, desc in features:
        tx = slide2.shapes.add_textbox(Inches(1.2), Inches(y_pos), Inches(10), Inches(0.9))
        tf = tx.text_frame
        tf.word_wrap = True
        p_title = tf.paragraphs[0]
        p_title.text = title
        r_title = p_title.runs[0]
        r_title.font.size = Pt(20)
        r_title.font.bold = True
        r_title.font.color.rgb = RGBColor(0x34, 0x78, 0xC6)

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        r_desc = p_desc.runs[0]
        r_desc.font.size = Pt(14)
        r_desc.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

        y_pos += 1.05

    # --- Slide 3: Technical Specifications Table ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    bg3 = slide3.background.fill
    bg3.solid()
    bg3.fore_color.rgb = RGBColor(0x0F, 0x1E, 0x30)

    tx_title3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(1))
    tf_t3 = tx_title3.text_frame
    p_t3 = tf_t3.paragraphs[0]
    p_t3.text = "Technical Specifications"
    r_t3 = p_t3.runs[0]
    r_t3.font.size = Pt(32)
    r_t3.font.bold = True
    r_t3.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    specs = [
        ("Dimensions", "324.9 x 226.8 x 14.9 mm"),
        ("Weight", "1.68 kg"),
        ("Display", '15.6" 3840x2160 OLED, 120Hz'),
        ("Battery", "84 Wh Li-Po, up to 18 hours"),
        ("Ports", "3x USB-C 4.0, 1x HDMI 2.1, SD card"),
        ("Operating System", "Linux / Windows 11 Pro"),
        ("Thermal Design", "Dual-fan vapor chamber, 45W TDP"),
        ("Audio", "Quad speakers, Dolby Atmos"),
    ]

    table_shape = slide3.shapes.add_table(len(specs) + 1, 2, Inches(1.5), Inches(1.5), Inches(10), Inches(5))
    table = table_shape.table
    table.columns[0].width = Inches(4)
    table.columns[1].width = Inches(6)

    # Header
    for ci, hdr in enumerate(["Parameter", "Value"]):
        cell = table.cell(0, ci)
        cell.text = hdr
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    for ri, (param, val) in enumerate(specs, 1):
        table.cell(ri, 0).text = param
        table.cell(ri, 1).text = val
        for ci in range(2):
            for run in table.cell(ri, ci).text_frame.paragraphs[0].runs:
                run.font.size = Pt(14)
                run.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

    # --- Slide 4: Product Image (the problematic slide) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])

    # Dark gradient-style background (solid dark since python-pptx gradient is limited)
    bg4 = slide4.background.fill
    bg4.solid()
    bg4.fore_color.rgb = RGBColor(0x1A, 0x0A, 0x2E)

    # Add a second dark overlay shape to simulate gradient effect
    overlay = slide4.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    overlay_fill = overlay.fill
    overlay_fill.solid()
    overlay_fill.fore_color.rgb = RGBColor(0x14, 0x08, 0x24)
    overlay.line.fill.background()  # no border

    tx_title4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(0.8))
    tf_t4 = tx_title4.text_frame
    p_t4 = tf_t4.paragraphs[0]
    p_t4.text = "Product Label & Certification"
    r_t4 = p_t4.runs[0]
    r_t4.font.size = Pt(28)
    r_t4.font.bold = True
    r_t4.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Add the transparent PNG image
    png_buf = create_transparent_png()
    img_left = Inches(4.0)
    img_top = Inches(1.8)
    img_width = Inches(5.0)
    img_height = Inches(3.75)
    pic = slide4.shapes.add_picture(png_buf, img_left, img_top, img_width, img_height)

    # Caption below image
    tx_cap = slide4.shapes.add_textbox(Inches(3.5), Inches(5.8), Inches(6), Inches(1.2))
    tf_cap = tx_cap.text_frame
    tf_cap.word_wrap = True
    p_cap = tf_cap.paragraphs[0]
    p_cap.text = "Note: The product label image above has transparency issues."
    p_cap.alignment = PP_ALIGN.CENTER
    r_cap = p_cap.runs[0]
    r_cap.font.size = Pt(14)
    r_cap.font.italic = True
    r_cap.font.color.rgb = RGBColor(0xFF, 0x99, 0x66)

    p_cap2 = tf_cap.add_paragraph()
    p_cap2.text = "The dark background shows through the transparent areas of the PNG."
    p_cap2.alignment = PP_ALIGN.CENTER
    r_cap2 = p_cap2.runs[0]
    r_cap2.font.size = Pt(12)
    r_cap2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    # --- Slide 5: Contact & Ordering ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    bg5 = slide5.background.fill
    bg5.solid()
    bg5.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)

    tx_title5 = slide5.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10), Inches(1))
    tf_t5 = tx_title5.text_frame
    p_t5 = tf_t5.paragraphs[0]
    p_t5.text = "Order & Contact Information"
    p_t5.alignment = PP_ALIGN.CENTER
    r_t5 = p_t5.runs[0]
    r_t5.font.size = Pt(36)
    r_t5.font.bold = True
    r_t5.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    contact_info = [
        "Sales: sales@technova-devices.com",
        "Support: support@technova-devices.com",
        "Phone: +1 (650) 555-0192",
        "Web: www.technova-devices.com/pro-x200",
        "Part Number: TNX200-PRO-64G-2T",
        "MSRP: $2,499.00 USD",
    ]
    tx_contact = slide5.shapes.add_textbox(Inches(3), Inches(3), Inches(7), Inches(3.5))
    tf_c = tx_contact.text_frame
    tf_c.word_wrap = True
    for i, line in enumerate(contact_info):
        if i == 0:
            p_c = tf_c.paragraphs[0]
        else:
            p_c = tf_c.add_paragraph()
        p_c.text = line
        r_c = p_c.runs[0]
        r_c.font.size = Pt(18)
        r_c.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
