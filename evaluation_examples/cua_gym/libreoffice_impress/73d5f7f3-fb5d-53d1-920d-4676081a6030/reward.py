"""
Reward Script: Add 3 textboxes to slide 2 arranged vertically with different text colors
Task ID: osworld_impress_textbox_colors_multiple_004
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4 pts): 3 textboxes with correct text content ('Item 1', 'Item 2', 'Item 3') exist on slide 2
  Component 2 (0.4 pts): Textboxes have correct colors (Item 1 = red, Item 2 = green, Item 3 = blue)
  Component 3 (0.2 pts): Textboxes are arranged vertically (tops are in strictly increasing order)
Total: 1.0
"""

import os

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError as e:
    print(f"CRITICAL: Cannot import python-pptx: {e}")
    print("REWARD: 0.0")
    exit(0)

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_textbox_colors_multiple_004'

# Expected text content and colors for each textbox
EXPECTED_ITEMS = [
    {'text': 'Item 1', 'color': 'FF0000', 'label': 'red'},
    {'text': 'Item 2', 'color': '008000', 'label': 'green'},
    {'text': 'Item 3', 'color': '0000FF', 'label': 'blue'},
]


def get_textboxes_on_slide2(prs):
    """Return all TEXT_BOX shapes on slide 2 (index 1), sorted by top position."""
    slide2 = prs.slides[1]
    textboxes = []
    for shape in slide2.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            textboxes.append(shape)
    # Sort by vertical position (top)
    textboxes.sort(key=lambda s: s.top if s.top is not None else 0)
    return textboxes


def get_shape_text(shape):
    """Extract all text from a shape's text frame as a stripped string."""
    if not shape.has_text_frame:
        return ''
    texts = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.text:
                texts.append(run.text)
    return ''.join(texts).strip()


def get_run_color(run):
    """Return hex RGB color string for a run, or None if not explicitly set."""
    try:
        if run.font.color.type is not None:
            return str(run.font.color.rgb).upper()
    except Exception:
        pass
    return None


def get_textbox_color(shape):
    """Return the font color of the first non-empty run in the textbox."""
    if not shape.has_text_frame:
        return None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if (run.text or '').strip():
                return get_run_color(run)
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has fewer than 2 slides (found {len(prs.slides)})")
        print("REWARD: 0.0")
        return 0.0

    # Retrieve textboxes on slide 2
    try:
        textboxes = get_textboxes_on_slide2(prs)
    except Exception as e:
        print(f"ERROR: Could not retrieve textboxes from slide 2: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: 3 textboxes with correct text content exist on slide 2 (0.4 points)
    # Each of 'Item 1', 'Item 2', 'Item 3' must be present as a textbox text
    try:
        found_texts = [get_shape_text(tb) for tb in textboxes]
        expected_texts = [item['text'] for item in EXPECTED_ITEMS]
        texts_matched = 0
        for expected_text in expected_texts:
            if expected_text in found_texts:
                texts_matched += 1

        if texts_matched == 3:
            print(f"PASS: Component 1 — All 3 textboxes found with correct text: {found_texts} (0.4 pts)")
            total_score += 0.4
        elif texts_matched > 0:
            partial = round(0.4 * texts_matched / 3, 4)
            print(f"PARTIAL: Component 1 — {texts_matched}/3 correct texts found: {found_texts} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — No correct textbox texts found. Found: {found_texts}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Textboxes have correct colors (Item 1=red FF0000, Item 2=green 008000, Item 3=blue 0000FF) (0.4 points)
    # Build a mapping from text to textbox for color checking
    try:
        text_to_tb = {}
        for tb in textboxes:
            t = get_shape_text(tb)
            text_to_tb[t] = tb

        colors_matched = 0
        for item in EXPECTED_ITEMS:
            tb = text_to_tb.get(item['text'])
            if tb is None:
                print(f"FAIL: Component 2 — Textbox with text '{item['text']}' not found for color check")
                continue
            actual_color = get_textbox_color(tb)
            if actual_color == item['color']:
                print(f"PASS: Component 2 — '{item['text']}' has correct color {item['label']} ({item['color']})")
                colors_matched += 1
            else:
                print(f"FAIL: Component 2 — '{item['text']}' expected color {item['color']} ({item['label']}), found {actual_color}")

        if colors_matched == 3:
            print(f"PASS: Component 2 — All 3 textboxes have correct colors (0.4 pts)")
            total_score += 0.4
        elif colors_matched > 0:
            partial = round(0.4 * colors_matched / 3, 4)
            print(f"PARTIAL: Component 2 — {colors_matched}/3 correct colors ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No textboxes have correct colors")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Textboxes are arranged vertically (tops are in strictly increasing order) (0.2 points)
    # This checks that the textboxes appear one above the other on the slide
    try:
        # Get sorted-by-top textboxes that match the expected texts (to avoid counting unrelated shapes)
        ordered_tbs = []
        for item in EXPECTED_ITEMS:
            tb = text_to_tb.get(item['text'])
            if tb is not None:
                ordered_tbs.append(tb)

        if len(ordered_tbs) == 3:
            tops = [tb.top for tb in ordered_tbs]
            # Verify tops are in strictly increasing order (vertically stacked)
            is_vertical = tops[0] < tops[1] < tops[2]
            if is_vertical:
                print(f"PASS: Component 3 — Textboxes are arranged vertically (tops: {[round(t/914400, 2) for t in tops]} inches) (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 3 — Textboxes are NOT arranged vertically. Tops (inches): {[round(t/914400, 2) for t in tops]}")
        else:
            print(f"FAIL: Component 3 — Could only find {len(ordered_tbs)} of 3 expected textboxes for vertical check")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {round(total_score, 4)}/1.0")
    print(f"REWARD: {round(final_score, 1)}")
    return final_score


# Default: test against canonical artifact path in a given env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
