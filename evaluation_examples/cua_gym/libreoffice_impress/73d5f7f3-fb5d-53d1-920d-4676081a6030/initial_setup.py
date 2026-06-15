"""
Initial Setup: 5-slide marketing deck presentation, slide 2 has only a title.
Task ID: osworld_impress_textbox_colors_multiple_004
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_textbox_colors_multiple_004'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard widescreen slide dimensions
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ----------------------------------------------------------------
    # Slide 1: Title Slide
    # ----------------------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Annual Marketing Strategy"
    slide1.placeholders[1].text = "Q1 2025 — Growth & Innovation"

    # ----------------------------------------------------------------
    # Slide 2: Section header — title ONLY (no textboxes; task adds them)
    # ----------------------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide2.shapes.title.text = "Key Product Lines"
    # Remove the content placeholder so slide 2 has only a title
    sp = slide2.placeholders[1]._element
    sp.getparent().remove(sp)

    # ----------------------------------------------------------------
    # Slide 3: Market Overview with bullet points
    # ----------------------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Market Overview"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Total Addressable Market: $4.2B"
    p3b = tf3.add_paragraph()
    p3b.text = "Year-over-year growth: 18%"
    p3c = tf3.add_paragraph()
    p3c.text = "Primary segments: Enterprise, SMB, Consumer"
    p3d = tf3.add_paragraph()
    p3d.text = "Geographic focus: North America, EMEA, APAC"

    # ----------------------------------------------------------------
    # Slide 4: Sales Performance
    # ----------------------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Sales Performance — Q1 2025"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Revenue: $12.8M (+22% YoY)"
    p4b = tf4.add_paragraph()
    p4b.text = "New customers acquired: 340"
    p4c = tf4.add_paragraph()
    p4c.text = "Average deal size: $37,600"
    p4d = tf4.add_paragraph()
    p4d.text = "Win rate: 43% (up from 38%)"
    p4e = tf4.add_paragraph()
    p4e.text = "Top region: North America — $7.1M"

    # ----------------------------------------------------------------
    # Slide 5: Next Steps & Call to Action
    # ----------------------------------------------------------------
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Next Steps"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Launch product refresh campaign — April 10"
    p5b = tf5.add_paragraph()
    p5b.text = "Onboard 3 new channel partners by May"
    p5c = tf5.add_paragraph()
    p5c.text = "Quarterly business reviews with top-20 accounts"
    p5d = tf5.add_paragraph()
    p5d.text = "Expand digital advertising budget by 15%"

    # ----------------------------------------------------------------
    # Save
    # ----------------------------------------------------------------
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
