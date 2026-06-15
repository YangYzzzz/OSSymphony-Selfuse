"""
Reward Script: Policy Brief Talking Points Import with KEY POINT Reordering
Task ID: osworld_multi_apps_impress_notes_import_012
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.3): All 10 slides have non-empty notes (notes were imported)
  - Component 2 (0.4): Slides 3, 6, 9 have KEY POINT lines moved to the beginning
  - Component 3 (0.3): Non-KEY POINT slides (1,2,4,5,7,8,10) have correct note content
"""

import os
from pptx import Presentation

WORKDIR = '/home/user/Desktop'
TARGET_FILE = f'{WORKDIR}/Policy_Brief_Final.pptx'

# Expected notes for each slide (1-indexed) based on KEY POINT reordering rules.
# Slides 3, 6, 9: KEY POINT lines moved to beginning.
# Others: as-is from the docx.
EXPECTED_NOTES = {
    1: [
        'Welcome members of the Senate Committee on Commerce, Science, and Transportation.',
        'This briefing represents 18 months of research from the National Infrastructure Foundation.',
        'We have consulted over 200 engineers, economists, and urban planners.',
        'Our findings are urgent and our recommendations are actionable.',
    ],
    2: [
        'The core argument is simple: deferred maintenance compounds costs.',
        'Each year of delay adds approximately $40 billion to the eventual repair bill.',
        'The window for cost-effective modernization is closing rapidly.',
        'Bipartisan support exists; what is needed is legislative action.',
    ],
    3: [
        # KEY POINT lines first
        'KEY POINT: Infrastructure grade has declined from D+ to D over the last decade.',
        'KEY POINT: Water system failures cause public health emergencies in 12 states annually.',
        # then non-KEY POINT lines
        'These statistics come from the American Society of Civil Engineers 2024 Report Card.',
        'The data on road conditions is from the Federal Highway Administration.',
        'Broadband access disparities disproportionately impact rural and tribal communities.',
    ],
    4: [
        'The $3,300 household cost includes vehicle repairs, commute delays, and utility bills.',
        'These are conservative estimates; actual costs may be 30% higher in rural areas.',
        'The job creation figures account for both construction and long-term operational roles.',
        'Return on investment data is supported by Congressional Budget Office modeling.',
    ],
    5: [
        'The high-speed rail proposal focuses on the Northeast and Pacific corridors first.',
        'Smart highway technology includes embedded sensors and vehicle-to-infrastructure communication.',
        'Port modernization is critical to supply chain resilience post-pandemic lessons.',
        'Urban transit funding prioritizes cities with populations over 250,000.',
    ],
    6: [
        # KEY POINT lines first
        'KEY POINT: Current grid infrastructure cannot support renewable targets without upgrade.',
        'KEY POINT: Grid failure risk increases 35% without transmission expansion by 2027.',
        # then non-KEY POINT lines
        'The 60% renewable target aligns with Paris Agreement commitments.',
        'Battery storage deployment will be distributed across 8 regional grid zones.',
        'Offshore wind zones have completed environmental impact assessments.',
        'Rural co-op modernization requires separate legislative authority under USDA programs.',
    ],
    7: [
        'Universal broadband was declared essential infrastructure by the FCC in 2021.',
        'The $65 billion rural figure represents a floor, not a ceiling.',
        'Cybersecurity integration must be built in from the start, not retrofitted.',
        '5G coverage on tribal lands requires consultation under Section 106 of NHPA.',
    ],
    8: [
        'The $400 billion in direct appropriations is spread over 5 annual budget cycles.',
        'Build America Bonds have a proven track record from the 2009 Recovery Act.',
        'Private sector commitments are contingent on regulatory certainty.',
        'Vehicle-miles-traveled fee pilots are already running in 8 states.',
    ],
    9: [
        # KEY POINT lines first
        'KEY POINT: Delaying Phase 1 by 6 months pushes the entire timeline back by 2 years.',
        'KEY POINT: Congressional appropriations must be multi-year to enable contractor certainty.',
        # then non-KEY POINT lines
        'Phase 1 funding can begin flowing within 90 days of enactment.',
        'Shovel-ready projects have already completed NEPA environmental review.',
        'High-speed rail construction in Phase 4 depends on right-of-way acquisition in Phases 1-2.',
        'Performance metrics will be publicly reported on a unified federal dashboard.',
    ],
    10: [
        'Recommendation 1 includes a 60-day fast-track committee review process.',
        'The National Infrastructure Council would include governors, mayors, and private sector leaders.',
        'One Federal Decision policy reduces permitting timelines from 4.5 years to 2 years.',
        'Workforce training will prioritize veterans, formerly incarcerated individuals, and rural workers.',
        'Thank you. We welcome questions from the committee.',
    ],
}

# Slides that require KEY POINT reordering (1-indexed)
KEY_POINT_SLIDES = {3, 6, 9}


def get_slide_notes(slide):
    """Get notes text from a slide, returning empty string if none."""
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def check_key_point_slide(slide_num, notes_text, expected_lines):
    """
    Check if KEY POINT lines are at the beginning of the notes for a given slide.
    Returns True if correctly placed, False otherwise.
    """
    notes_lines = [line for line in notes_text.split('\n') if line.strip()]
    expected_key_points = [l for l in expected_lines if l.startswith('KEY POINT:')]
    n_kp = len(expected_key_points)
    actual_first_n = notes_lines[:n_kp]

    # Verify each of the first N lines is a KEY POINT line
    for j, (actual_line, expected_kp) in enumerate(zip(actual_first_n, expected_key_points)):
        # Check if the line starts with KEY POINT: and has matching content
        if not actual_line.startswith('KEY POINT:'):
            print(f"FAIL: Component 2 — Slide {slide_num} line {j+1}: expected KEY POINT:, got: {repr(actual_line[:80])}")
            return False
        # Also verify rough content match
        key_phrase = expected_kp[11:50].lower()  # skip 'KEY POINT: '
        if key_phrase not in actual_line.lower():
            print(f"FAIL: Component 2 — Slide {slide_num} line {j+1}: KEY POINT content mismatch. Expected: {repr(expected_kp[:60])}, got: {repr(actual_line[:60])}")
            return False

    if len(actual_first_n) < n_kp:
        print(f"FAIL: Component 2 — Slide {slide_num}: expected {n_kp} KEY POINT lines first, found only {len(actual_first_n)}")
        return False

    # Verify no KEY POINT lines remain after the first section
    remaining_lines = notes_lines[n_kp:]
    kp_in_remaining = [l for l in remaining_lines if l.startswith('KEY POINT:')]
    if kp_in_remaining:
        print(f"FAIL: Component 2 — Slide {slide_num}: KEY POINT lines found after first section: {kp_in_remaining}")
        return False

    return True


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Gate: Must have 10 slides
    if len(prs.slides) != 10:
        print(f"CRITICAL: Expected 10 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Collect all slide notes
    slide_notes = {}
    for i, slide in enumerate(prs.slides):
        slide_notes[i + 1] = get_slide_notes(slide)

    # Component 1: All 10 slides have non-empty notes (0.3 points)
    # This tests that notes were imported at all (fails on initial where notes are empty)
    try:
        empty_slides = [i for i in range(1, 11) if not slide_notes[i].strip()]
        if not empty_slides:
            print("PASS: Component 1 — All 10 slides have non-empty notes (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 — {len(empty_slides)} slides have empty notes: {empty_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slides 3, 6, 9 have KEY POINT lines at the beginning (0.4 points)
    # Each KEY POINT slide contributes equal share
    try:
        key_point_slides_sorted = sorted(KEY_POINT_SLIDES)
        per_slide_pts = round(0.4 / len(key_point_slides_sorted), 4)

        for slide_num in key_point_slides_sorted:
            notes_text = slide_notes[slide_num]
            expected_lines = EXPECTED_NOTES[slide_num]

            if not notes_text.strip():
                print(f"FAIL: Component 2 — Slide {slide_num} has no notes text")
                continue

            if check_key_point_slide(slide_num, notes_text, expected_lines):
                print(f"PASS: Component 2 — Slide {slide_num} KEY POINT lines correctly placed first ({per_slide_pts} pts)")
                total_score += per_slide_pts
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Non-KEY POINT slides (1,2,4,5,7,8,10) have correct note content (0.3 points)
    # Each correct slide contributes partial credit
    try:
        non_kp_slides = [s for s in range(1, 11) if s not in KEY_POINT_SLIDES]
        per_slide_pts_3 = round(0.3 / len(non_kp_slides), 4)

        for slide_num in non_kp_slides:
            notes_text = slide_notes[slide_num]
            expected_lines = EXPECTED_NOTES[slide_num]

            if not notes_text.strip():
                print(f"FAIL: Component 3 — Slide {slide_num} has empty notes")
                continue

            # Check that all expected lines are present in the notes
            notes_lower = notes_text.lower()
            missing_lines = []
            for expected_line in expected_lines:
                # Use a key phrase check (first 40 chars) to avoid whitespace/encoding issues
                key_phrase = expected_line[:40].lower()
                if key_phrase not in notes_lower:
                    missing_lines.append(expected_line[:60])

            if not missing_lines:
                print(f"PASS: Component 3 — Slide {slide_num} has correct notes content ({per_slide_pts_3} pts)")
                total_score += per_slide_pts_3
            else:
                print(f"FAIL: Component 3 — Slide {slide_num} missing lines: {missing_lines}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 4)
    print(f"\nScore: {total_score:.4f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Check if the output file exists
if not os.path.exists(TARGET_FILE):
    print(f"File not found: {TARGET_FILE}")
    print("REWARD: 0.0")
else:
    verify_task(TARGET_FILE)
