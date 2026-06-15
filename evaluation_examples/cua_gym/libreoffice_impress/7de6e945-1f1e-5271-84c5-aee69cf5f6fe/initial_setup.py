"""
Initial Setup: Policy briefing presentation and talking points document
Task ID: osworld_multi_apps_impress_notes_import_012
Domain: libreoffice_impress

Creates:
  - /home/user/Desktop/Policy_Brief.pptx  (10 slides, empty notes)
  - /home/user/Desktop/policy_talking_points.docx  (notes for all 10 slides, 3 with KEY POINT: lines)
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Pt as DocPt

DESKTOP = '/home/user/Desktop'
PPTX_PATH = f'{DESKTOP}/Policy_Brief.pptx'
DOCX_PATH = f'{DESKTOP}/policy_talking_points.docx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_presentation():
    """Create Policy_Brief.pptx with 10 slides and empty notes."""
    prs = Presentation()

    slide_data = [
        {
            "title": "National Infrastructure Modernization Policy Brief",
            "subtitle": "Federal Transportation and Energy Strategy 2025-2030",
            "layout": 0,
        },
        {
            "title": "Executive Summary",
            "body": (
                "The United States faces a critical juncture in infrastructure investment.\n"
                "Aging transportation networks and energy grids require immediate modernization.\n"
                "This brief outlines key recommendations for Congressional action.\n"
                "Estimated investment: $1.2 trillion over 5 years."
            ),
            "layout": 1,
        },
        {
            "title": "Current State of U.S. Infrastructure",
            "body": (
                "Roads & Bridges: 43% of public roads in poor or mediocre condition\n"
                "Rail Networks: Average passenger rail speed 50% below European peers\n"
                "Energy Grid: 70% of transmission lines over 25 years old\n"
                "Broadband: 21 million Americans lack reliable internet access\n"
                "Water Systems: 6 billion gallons lost daily to leaky pipes"
            ),
            "layout": 1,
        },
        {
            "title": "Economic Impact Analysis",
            "body": (
                "Infrastructure deficiencies cost U.S. households $3,300 per year.\n"
                "Delayed freight shipments reduce GDP by an estimated 0.4% annually.\n"
                "Energy grid inefficiencies cost businesses $150 billion each year.\n"
                "Job creation potential: 11 million direct and indirect jobs.\n"
                "Every $1 invested returns $2.20 in long-term economic value."
            ),
            "layout": 1,
        },
        {
            "title": "Transportation Modernization",
            "body": (
                "High-Speed Rail Corridors: Connect 20 major metros by 2028\n"
                "Smart Highway Initiative: Autonomous vehicle infrastructure on I-95 and I-10\n"
                "Port Modernization: Upgrade top 15 ports for next-generation cargo vessels\n"
                "Urban Transit: $200 billion for metropolitan mass transit expansion\n"
                "Aviation Overhaul: NextGen air traffic control system nationwide"
            ),
            "layout": 1,
        },
        {
            "title": "Clean Energy Grid Transition",
            "body": (
                "Renewable Energy Target: 60% clean energy by 2030\n"
                "Grid Storage: 100 GWh battery storage deployment\n"
                "Transmission Expansion: 40,000 miles of new high-voltage lines\n"
                "Offshore Wind: 30 GW capacity in Atlantic and Pacific coastal zones\n"
                "Rural Electrification: Modernize co-op grids serving 56 million Americans"
            ),
            "layout": 1,
        },
        {
            "title": "Digital Infrastructure",
            "body": (
                "National Broadband Plan: Universal high-speed access by 2027\n"
                "Rural Connectivity: $65 billion targeted at underserved communities\n"
                "Cybersecurity Framework: Unified protection for critical infrastructure\n"
                "5G Deployment: Full nationwide coverage including tribal lands\n"
                "Open Access Policy: Prevent monopolization of last-mile networks"
            ),
            "layout": 1,
        },
        {
            "title": "Funding Mechanisms",
            "body": (
                "Federal Appropriations: $400 billion over 5 years from general revenue\n"
                "Infrastructure Bonds: $300 billion in Build America Bonds\n"
                "Public-Private Partnerships: $250 billion leveraged from private sector\n"
                "Gas Tax Modernization: Indexed to inflation + vehicle-miles-traveled fee\n"
                "State Matching Grants: 3:1 federal match for qualified state projects"
            ),
            "layout": 1,
        },
        {
            "title": "Implementation Timeline",
            "body": (
                "Phase 1 (2025-2026): Emergency bridge repairs + shovel-ready projects\n"
                "Phase 2 (2026-2027): Begin major corridor construction\n"
                "Phase 3 (2027-2028): Energy grid upgrades + rural broadband rollout\n"
                "Phase 4 (2028-2029): High-speed rail construction accelerates\n"
                "Phase 5 (2029-2030): Final integration and performance evaluation"
            ),
            "layout": 1,
        },
        {
            "title": "Recommendations & Next Steps",
            "body": (
                "Recommendation 1: Pass Infrastructure Investment Act by Q2 2025\n"
                "Recommendation 2: Establish bipartisan National Infrastructure Council\n"
                "Recommendation 3: Streamline permitting via One Federal Decision policy\n"
                "Recommendation 4: Create Infrastructure Workforce Training Program\n"
                "Recommendation 5: Launch annual infrastructure performance report"
            ),
            "layout": 1,
        },
    ]

    for i, data in enumerate(slide_data):
        layout_idx = data["layout"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = data["title"]

        # Set body/subtitle content
        if layout_idx == 0 and "subtitle" in data:
            # Title slide
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = data["subtitle"]
        elif layout_idx == 1 and "body" in data:
            # Title + content
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = data["body"]

        # Notes are intentionally left EMPTY for initial state

    prs.save(PPTX_PATH)
    print(f'Presentation created: {PPTX_PATH}')


def create_talking_points():
    """
    Create policy_talking_points.docx with notes for all 10 slides.
    Slides 3, 6, and 9 (1-indexed) have KEY POINT: lines mixed within their text.
    """
    doc = Document()
    doc.add_heading('Policy Brief Talking Points', level=0)

    # Slide-by-slide notes
    # KEY POINT: lines appear mid-text in slides 3, 6, and 9
    slides_notes = [
        # Slide 1
        (
            "Slide 1",
            [
                "Welcome members of the Senate Committee on Commerce, Science, and Transportation.",
                "This briefing represents 18 months of research from the National Infrastructure Foundation.",
                "We have consulted over 200 engineers, economists, and urban planners.",
                "Our findings are urgent and our recommendations are actionable.",
            ]
        ),
        # Slide 2
        (
            "Slide 2",
            [
                "The core argument is simple: deferred maintenance compounds costs.",
                "Each year of delay adds approximately $40 billion to the eventual repair bill.",
                "The window for cost-effective modernization is closing rapidly.",
                "Bipartisan support exists; what is needed is legislative action.",
            ]
        ),
        # Slide 3 — contains KEY POINT: lines
        (
            "Slide 3",
            [
                "These statistics come from the American Society of Civil Engineers 2024 Report Card.",
                "KEY POINT: Infrastructure grade has declined from D+ to D over the last decade.",
                "The data on road conditions is from the Federal Highway Administration.",
                "Broadband access disparities disproportionately impact rural and tribal communities.",
                "KEY POINT: Water system failures cause public health emergencies in 12 states annually.",
            ]
        ),
        # Slide 4
        (
            "Slide 4",
            [
                "The $3,300 household cost includes vehicle repairs, commute delays, and utility bills.",
                "These are conservative estimates; actual costs may be 30% higher in rural areas.",
                "The job creation figures account for both construction and long-term operational roles.",
                "Return on investment data is supported by Congressional Budget Office modeling.",
            ]
        ),
        # Slide 5
        (
            "Slide 5",
            [
                "The high-speed rail proposal focuses on the Northeast and Pacific corridors first.",
                "Smart highway technology includes embedded sensors and vehicle-to-infrastructure communication.",
                "Port modernization is critical to supply chain resilience post-pandemic lessons.",
                "Urban transit funding prioritizes cities with populations over 250,000.",
            ]
        ),
        # Slide 6 — contains KEY POINT: lines
        (
            "Slide 6",
            [
                "The 60% renewable target aligns with Paris Agreement commitments.",
                "Battery storage deployment will be distributed across 8 regional grid zones.",
                "KEY POINT: Current grid infrastructure cannot support renewable targets without upgrade.",
                "Offshore wind zones have completed environmental impact assessments.",
                "Rural co-op modernization requires separate legislative authority under USDA programs.",
                "KEY POINT: Grid failure risk increases 35% without transmission expansion by 2027.",
            ]
        ),
        # Slide 7
        (
            "Slide 7",
            [
                "Universal broadband was declared essential infrastructure by the FCC in 2021.",
                "The $65 billion rural figure represents a floor, not a ceiling.",
                "Cybersecurity integration must be built in from the start, not retrofitted.",
                "5G coverage on tribal lands requires consultation under Section 106 of NHPA.",
            ]
        ),
        # Slide 8
        (
            "Slide 8",
            [
                "The $400 billion in direct appropriations is spread over 5 annual budget cycles.",
                "Build America Bonds have a proven track record from the 2009 Recovery Act.",
                "Private sector commitments are contingent on regulatory certainty.",
                "Vehicle-miles-traveled fee pilots are already running in 8 states.",
            ]
        ),
        # Slide 9 — contains KEY POINT: lines
        (
            "Slide 9",
            [
                "Phase 1 funding can begin flowing within 90 days of enactment.",
                "Shovel-ready projects have already completed NEPA environmental review.",
                "KEY POINT: Delaying Phase 1 by 6 months pushes the entire timeline back by 2 years.",
                "High-speed rail construction in Phase 4 depends on right-of-way acquisition in Phases 1-2.",
                "Performance metrics will be publicly reported on a unified federal dashboard.",
                "KEY POINT: Congressional appropriations must be multi-year to enable contractor certainty.",
            ]
        ),
        # Slide 10
        (
            "Slide 10",
            [
                "Recommendation 1 includes a 60-day fast-track committee review process.",
                "The National Infrastructure Council would include governors, mayors, and private sector leaders.",
                "One Federal Decision policy reduces permitting timelines from 4.5 years to 2 years.",
                "Workforce training will prioritize veterans, formerly incarcerated individuals, and rural workers.",
                "Thank you. We welcome questions from the committee.",
            ]
        ),
    ]

    for slide_label, notes_lines in slides_notes:
        doc.add_heading(slide_label, level=1)
        for line in notes_lines:
            para = doc.add_paragraph(line)

    doc.save(DOCX_PATH)
    print(f'Talking points document created: {DOCX_PATH}')


def main():
    os.makedirs(DESKTOP, exist_ok=True)

    create_presentation()
    create_talking_points()

    # GUI-ready startup: open Policy_Brief.pptx in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{PPTX_PATH}"', delay_sec=3.0)
    print('GUI_READY: Launched LibreOffice Impress with Policy_Brief.pptx (DISPLAY=:0)')


main()
