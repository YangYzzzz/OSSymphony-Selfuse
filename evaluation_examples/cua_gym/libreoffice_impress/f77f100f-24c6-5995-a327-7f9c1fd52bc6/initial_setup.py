"""
Initial Setup: Create a 12-slide research talk presentation (no section dividers)
Task ID: impress_stu_068
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_068'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, alignment=PP_ALIGN.LEFT, color=RGBColor(0x33, 0x33, 0x33),
                 font_name="Arial"):
    """Helper to add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=RGBColor(0x33, 0x33, 0x33)):
    """Helper to add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    # Standard widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # ============================================================
    # Slide 1: Title Slide
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(2),
                 "Adaptive Neural Network Architectures\nfor Climate Prediction Models",
                 font_size=36, bold=True, alignment=PP_ALIGN.CENTER,
                 color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(1.5),
                 "Dr. Elena Vasquez, Prof. James Whitfield\n"
                 "Institute for Computational Earth Sciences\nMarch 2025",
                 font_size=18, alignment=PP_ALIGN.CENTER,
                 color=RGBColor(0xCC, 0xCC, 0xCC))

    # ============================================================
    # Slide 2: Agenda / Overview
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "Agenda", font_size=32, bold=True,
                 color=RGBColor(0x1B, 0x3A, 0x5C))
    add_bullet_list(slide, Inches(1), Inches(1.8), Inches(10), Inches(4.5), [
        "1. Background and motivation",
        "2. Theoretical framework for neural climate models",
        "3. Proposed architecture: ClimateNet-V3",
        "4. Experimental methodology and data pipelines",
        "5. Key results and comparative analysis",
        "6. Ablation studies and sensitivity analysis",
        "7. Conclusions and future directions",
    ], font_size=20)

    # ============================================================
    # Slide 3: Background - Climate Modeling Challenges
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "Climate Modeling Challenges", font_size=28, bold=True,
                 color=RGBColor(0x1B, 0x3A, 0x5C))
    add_bullet_list(slide, Inches(1), Inches(1.8), Inches(5.5), Inches(4.5), [
        "Traditional GCMs require enormous computational resources",
        "Resolution limitations: 100km grid cells miss local phenomena",
        "Parameterization of sub-grid processes introduces systematic biases",
        "Ensemble runs for uncertainty quantification are prohibitively expensive",
        "Data assimilation from heterogeneous sensor networks remains difficult",
    ], font_size=16)
    add_text_box(slide, Inches(7), Inches(2), Inches(5.5), Inches(4),
                 "Key Statistic:\nGlobal climate simulations at 1km resolution "
                 "would require 10^18 FLOPS per simulated year, exceeding current "
                 "exascale computing capacity by two orders of magnitude.",
                 font_size=14, color=RGBColor(0x66, 0x66, 0x66))

    # ============================================================
    # Slide 4: Literature Review
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "Prior Work in Neural Climate Prediction", font_size=28, bold=True,
                 color=RGBColor(0x1B, 0x3A, 0x5C))
    add_bullet_list(slide, Inches(1), Inches(1.8), Inches(10), Inches(5), [
        "Weyn et al. (2020): CNN-based weather forecasting with cubed-sphere mapping",
        "Rasp & Thuerey (2021): Data-driven medium-range weather prediction",
        "Pathak et al. (2022): FourCastNet using adaptive Fourier neural operators",
        "Bi et al. (2023): Pangu-Weather with 3D Earth-specific transformer",
        "Lam et al. (2023): GraphCast achieving SOTA 10-day forecasts",
        "Limitation: Most focus on short-term weather, not decadal climate projections",
    ], font_size=16)

    # ============================================================
    # Slide 5: Theoretical Framework
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "Theoretical Framework", font_size=28, bold=True,
                 color=RGBColor(0x1B, 0x3A, 0x5C))
    add_text_box(slide, Inches(1), Inches(1.8), Inches(10), Inches(1.5),
                 "We formulate climate prediction as a sequence-to-sequence problem "
                 "over spatiotemporal fields, incorporating physical conservation laws "
                 "as differentiable constraints in the loss function.",
                 font_size=18, color=RGBColor(0x33, 0x33, 0x33))
    add_bullet_list(slide, Inches(1), Inches(3.5), Inches(10), Inches(3), [
        "Conservation of energy: enforced via Lagrangian penalty term",
        "Mass continuity: embedded in architecture through divergence-free layers",
        "Radiative transfer: approximated with physics-informed neural operator blocks",
        "Temporal coherence: multi-scale attention over 1-day, 1-month, and 1-year windows",
    ], font_size=16)

    # ============================================================
    # Slide 6: ClimateNet-V3 Architecture
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "ClimateNet-V3 Architecture", font_size=28, bold=True,
                 color=RGBColor(0x1B, 0x3A, 0x5C))
    add_bullet_list(slide, Inches(1), Inches(1.8), Inches(5), Inches(4.5), [
        "Encoder: Spherical harmonic embedding + Vision Transformer",
        "Processor: 24-layer graph neural network on icosahedral mesh",
        "Decoder: Adaptive upsampling with skip connections to encoder",
        "Parameters: 1.2 billion (760M processor, 280M encoder, 160M decoder)",
        "Training: Mixed-precision on 256 A100 GPUs for 14 days",
    ], font_size=16)
    add_text_box(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(5),
                 "Model Specifications:\n\n"
                 "Input resolution: 0.25 degree (~28km)\n"
                 "Output resolution: 0.25 degree\n"
                 "Temporal step: 6 hours\n"
                 "Variables: 73 atmospheric + 12 surface\n"
                 "Pressure levels: 13 (1000-50 hPa)\n"
                 "Memory footprint: 18.4 GB",
                 font_size=14, color=RGBColor(0x44, 0x44, 0x44))

    # ============================================================
    # Slide 7: Training Data and Pipeline
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "Training Data and Pipeline", font_size=28, bold=True,
                 color=RGBColor(0x1B, 0x3A, 0x5C))
    add_bullet_list(slide, Inches(1), Inches(1.8), Inches(10), Inches(5), [
        "ERA5 reanalysis: 1979-2022 (43 years, 6-hourly, ~2.8 PB total)",
        "CMIP6 ensemble: 35 models, 4 SSP scenarios, 2015-2100",
        "Satellite observations: MODIS, AIRS, CERES for validation",
        "Data augmentation: temporal jittering, hemispheric flipping",
        "Train/Val/Test split: 1979-2015 / 2016-2019 / 2020-2022",
        "Preprocessing: z-score normalization per variable per pressure level",
    ], font_size=16)

    # ============================================================
    # Slide 8: Experimental Setup
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "Experimental Setup", font_size=28, bold=True,
                 color=RGBColor(0x1B, 0x3A, 0x5C))
    add_bullet_list(slide, Inches(1), Inches(1.8), Inches(10), Inches(5), [
        "Baselines: IFS-HRES (ECMWF), GraphCast, Pangu-Weather, FourCastNet",
        "Metrics: RMSE, ACC (anomaly correlation coefficient), CRPS for probabilistic",
        "Evaluation lead times: 1, 3, 5, 7, 10, 14, 30, 90, 365 days",
        "Regional breakdowns: tropics, midlatitudes, polar regions",
        "Extreme event detection: tropical cyclone tracking, heatwave identification",
        "Computational cost comparison: FLOPS per forecast step",
    ], font_size=16)

    # ============================================================
    # Slide 9: Key Results - Temperature Forecasting
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "Results: Temperature Forecasting (500hPa)", font_size=28, bold=True,
                 color=RGBColor(0x1B, 0x3A, 0x5C))
    # Create a results table
    table_shape = slide.shapes.add_table(5, 4, Inches(1.5), Inches(2), Inches(9), Inches(3.5))
    table = table_shape.table
    headers = ["Model", "3-Day RMSE (K)", "10-Day RMSE (K)", "ACC @ 10d"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    data = [
        ["IFS-HRES", "1.42", "3.87", "0.891"],
        ["GraphCast", "1.38", "3.62", "0.904"],
        ["Pangu-Weather", "1.45", "3.71", "0.897"],
        ["ClimateNet-V3", "1.21", "3.18", "0.932"],
    ]
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val
            for run in table.cell(r, c).text_frame.paragraphs[0].runs:
                run.font.size = Pt(13)

    # ============================================================
    # Slide 10: Ablation Study
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "Ablation Study", font_size=28, bold=True,
                 color=RGBColor(0x1B, 0x3A, 0x5C))
    add_bullet_list(slide, Inches(1), Inches(1.8), Inches(10), Inches(5), [
        "Without physics constraints: RMSE increases 12% at 10-day lead time",
        "Without multi-scale attention: ACC drops from 0.932 to 0.908",
        "Without spherical harmonic embedding: polar region errors increase 23%",
        "Reduced processor layers (12 vs 24): 7% degradation across all metrics",
        "Standard positional encoding vs spherical: 15% worse in tropics",
    ], font_size=16)

    # ============================================================
    # Slide 11: Discussion and Limitations
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "Discussion and Limitations", font_size=28, bold=True,
                 color=RGBColor(0x1B, 0x3A, 0x5C))
    add_bullet_list(slide, Inches(1), Inches(1.8), Inches(10), Inches(5), [
        "ClimateNet-V3 achieves SOTA on medium-range forecasting benchmarks",
        "Physics constraints improve long-range stability significantly",
        "Limitation: decadal projections still show drift after 5+ years",
        "Limitation: extreme precipitation events remain challenging",
        "Training cost remains high: ~$180K for a single full training run",
        "Open question: transferability to paleoclimate reconstruction tasks",
    ], font_size=16)

    # ============================================================
    # Slide 12: Future Work and Acknowledgments
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    add_text_box(slide, Inches(1), Inches(1), Inches(11), Inches(1.5),
                 "Future Directions & Acknowledgments",
                 font_size=32, bold=True, alignment=PP_ALIGN.CENTER,
                 color=RGBColor(0xFF, 0xFF, 0xFF))
    add_bullet_list(slide, Inches(1.5), Inches(3), Inches(9.5), Inches(3.5), [
        "Extend to coupled ocean-atmosphere modeling",
        "Integrate with carbon cycle and ice sheet dynamics",
        "Develop efficient fine-tuning for regional downscaling",
        "Release pre-trained weights and evaluation benchmarks",
        "",
        "Funded by NSF Grant #2345678 and DOE ASCR Program",
        "Computing resources: NERSC Perlmutter supercomputer",
    ], font_size=18, color=RGBColor(0xDD, 0xDD, 0xDD))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
