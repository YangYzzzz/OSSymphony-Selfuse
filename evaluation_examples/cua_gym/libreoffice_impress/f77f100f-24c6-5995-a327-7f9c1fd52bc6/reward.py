"""
Reward Script: Section Divider System for Presentation
Task ID: impress_stu_068
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): Total slide count is 15 (3 dividers inserted)
  Component 2 (0.25): Divider slides have teal (#008080) background
  Component 3 (0.25): Divider slides have correct section name text, bold, white, ~44pt
  Component 4 (0.15): Divider slides have two white horizontal lines each
  Component 5 (0.15): Dissolve transition on divider slides only
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_068'

# Expected divider positions (0-indexed): slides 3, 7, 11 in the 15-slide presentation
DIVIDER_INDICES = [2, 6, 10]
SECTION_NAMES = ['Part I: Theory', 'Part II: Methods', 'Part III: Results']
EXPECTED_TOTAL_SLIDES = 15


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Slide count is 15 (0.20 points)
    try:
        if num_slides == EXPECTED_TOTAL_SLIDES:
            print(f"PASS: Component 1 -- Slide count is {num_slides} (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 -- Expected {EXPECTED_TOTAL_SLIDES} slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # For remaining components, we need to identify divider slides.
    # We check slides at expected positions if total count matches,
    # otherwise search for divider-like slides.
    divider_slides = []
    divider_section_map = {}
    if num_slides == EXPECTED_TOTAL_SLIDES:
        for idx, expected_name in zip(DIVIDER_INDICES, SECTION_NAMES):
            divider_slides.append(idx)
            divider_section_map[idx] = expected_name
    else:
        # Fallback: scan for slides with #008080 background
        for idx in range(num_slides):
            slide = prs.slides[idx]
            try:
                fill = slide.background.fill
                if fill.type == 1:  # SOLID
                    rgb_str = str(fill.fore_color.rgb)
                    if rgb_str == '008080':
                        divider_slides.append(idx)
            except Exception:
                pass
        # Map found dividers to expected names by order
        for i, idx in enumerate(divider_slides[:3]):
            divider_section_map[idx] = SECTION_NAMES[i] if i < len(SECTION_NAMES) else ''

    # Component 2: Divider slides have teal (#008080) background (0.25 points)
    try:
        bg_pass_count = 0
        for idx in divider_slides[:3]:
            slide = prs.slides[idx]
            fill = slide.background.fill
            if fill.type == 1:  # SOLID
                rgb_str = str(fill.fore_color.rgb)
                if rgb_str == '008080':
                    bg_pass_count += 1
                else:
                    print(f"FAIL: Component 2 -- Slide {idx+1} background is #{rgb_str}, expected #008080")
            else:
                print(f"FAIL: Component 2 -- Slide {idx+1} background fill type is {fill.type}, expected SOLID (1)")

        if bg_pass_count == 3:
            print(f"PASS: Component 2 -- All 3 divider slides have #008080 teal background (0.25 pts)")
            total_score += 0.25
        elif bg_pass_count > 0:
            partial = round(0.25 * bg_pass_count / 3, 2)
            print(f"PARTIAL: Component 2 -- {bg_pass_count}/3 dividers have correct background ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 -- No divider slides have #008080 background")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Section name text, bold, white, ~44pt (0.25 points)
    try:
        text_pass_count = 0
        for idx in divider_slides[:3]:
            slide = prs.slides[idx]
            expected_name = divider_section_map.get(idx, '')
            found_text = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        para_text = para.text.strip()
                        if expected_name and expected_name in para_text:
                            # Check runs for formatting
                            runs = [r for r in para.runs if (r.text or '').strip()]
                            if not runs:
                                print(f"FAIL: Component 3 -- Slide {idx+1} has text '{para_text}' but no formatted runs")
                                continue
                            run = runs[0]
                            # Check bold
                            is_bold = run.font.bold is True
                            # Check color white
                            try:
                                color_str = str(run.font.color.rgb)
                                is_white = color_str == 'FFFFFF'
                            except Exception:
                                is_white = False
                            # Check size ~44pt (558800 EMU, allow tolerance)
                            font_size = run.font.size
                            is_size_ok = font_size is not None and abs(font_size - 558800) < 25400  # ~2pt tolerance

                            if is_bold and is_white and is_size_ok:
                                text_pass_count += 1
                                found_text = True
                            else:
                                details = f"bold={is_bold}, white={is_white}, size={font_size}"
                                print(f"FAIL: Component 3 -- Slide {idx+1} text '{para_text}' formatting wrong: {details}")
                                found_text = True
                            break
                if found_text:
                    break
            if not found_text:
                print(f"FAIL: Component 3 -- Slide {idx+1} missing expected text '{expected_name}'")

        if text_pass_count == 3:
            print(f"PASS: Component 3 -- All 3 dividers have correct section names with formatting (0.25 pts)")
            total_score += 0.25
        elif text_pass_count > 0:
            partial = round(0.25 * text_pass_count / 3, 2)
            print(f"PARTIAL: Component 3 -- {text_pass_count}/3 dividers have correct text ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 -- No divider slides have correct section names")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Two white horizontal lines on each divider slide (0.15 points)
    try:
        lines_pass_count = 0
        for idx in divider_slides[:3]:
            slide = prs.slides[idx]
            white_lines = []
            for shape in slide.shapes:
                if shape.shape_type == 9:  # LINE / CONNECTOR
                    try:
                        line_rgb = str(shape.line.color.rgb)
                        if line_rgb == 'FFFFFF':
                            white_lines.append(shape)
                    except Exception:
                        pass
            if len(white_lines) >= 2:
                lines_pass_count += 1
            else:
                print(f"FAIL: Component 4 -- Slide {idx+1} has {len(white_lines)} white lines, expected 2")

        if lines_pass_count == 3:
            print(f"PASS: Component 4 -- All 3 dividers have 2 white horizontal lines (0.15 pts)")
            total_score += 0.15
        elif lines_pass_count > 0:
            partial = round(0.15 * lines_pass_count / 3, 2)
            print(f"PARTIAL: Component 4 -- {lines_pass_count}/3 dividers have correct lines ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 -- No divider slides have 2 white horizontal lines")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Dissolve transition on divider slides only (0.15 points)
    try:
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        divider_trans_ok = 0
        non_divider_trans_clean = True

        with zipfile.ZipFile(file_path, 'r') as zf:
            for slide_idx in range(num_slides):
                fname = f'ppt/slides/slide{slide_idx + 1}.xml'
                try:
                    with zf.open(fname) as f:
                        root = ET.parse(f).getroot()
                        tr = root.find('.//p:transition', ns)
                        has_dissolve = (tr is not None and tr.find('.//p:dissolve', ns) is not None)

                        if slide_idx in divider_slides[:3]:
                            if has_dissolve:
                                divider_trans_ok += 1
                            else:
                                print(f"FAIL: Component 5 -- Divider slide {slide_idx+1} missing Dissolve transition")
                        else:
                            if has_dissolve:
                                non_divider_trans_clean = False
                                print(f"FAIL: Component 5 -- Non-divider slide {slide_idx+1} has Dissolve transition (should not)")
                except KeyError:
                    pass

        if divider_trans_ok == 3 and non_divider_trans_clean:
            print(f"PASS: Component 5 -- Dissolve transition on all 3 dividers only (0.15 pts)")
            total_score += 0.15
        elif divider_trans_ok > 0:
            # Partial: credit for dividers having transition, small penalty for non-divider having it
            partial = round(0.15 * divider_trans_ok / 3, 2)
            if not non_divider_trans_clean:
                partial = round(partial * 0.5, 2)  # penalize for extra transitions
            print(f"PARTIAL: Component 5 -- {divider_trans_ok}/3 dividers have transition ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 -- No divider slides have Dissolve transition")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
