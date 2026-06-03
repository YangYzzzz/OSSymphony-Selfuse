"""
FINAL REWARD SCRIPT - SUCCESS
Task: While rehearsing, I noticed slide 245 doesn’t let me move around the deck smoothly. In LibreOffice Impress, how do I turn the logo on that slide into a hyperlink that jumps straight to slide 1, and make the text box that literally says "Next" link directly to slide 2?
Generated: 2025-09-10 19:05:10
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
import re
from pptx import Presentation

def verify_hyperlinks(file_path: str) -> float:
    """Verify that on slide 245 the logo links to slide 1 and the text box
    containing the word "Next" links to slide 2.

    Scoring:
        +0.5  Logo correctly linked to slide 1 (index 0)
        +0.5  "Next" text box correctly linked to slide 2 (index 1)
    Returns a float between 0.0 and 1.0
    """

    print("Starting verification of internal hyperlinks …")
    score = 0.0

    # ---------- 1. Load presentation safely ---------- #
    if not os.path.exists(file_path):
        print("✗ File not found:", file_path)
        return 0.0

    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded successfully with {len(prs.slides)} slides")
    except Exception as exc:
        print("✗ Failed to load PPTX:", exc)
        return 0.0

    # ---------- 2. Ensure slide 245 exists ---------- #
    if len(prs.slides) < 245:
        print("✗ Presentation has fewer than 245 slides – cannot verify task requirements")
        return 0.0

    slide_245 = prs.slides[244]  # zero-based index

    # Build helper dict: slide_id → index for quick lookup
    slide_id_to_index = {s.slide_id: idx for idx, s in enumerate(prs.slides)}

    # Flags we are looking for
    logo_link_found = False  # hyperlink to slide 1 (index 0)
    next_link_found = False  # text containing "Next" linking to slide 2 (index 1)

    # ---------- 3. Inspect every shape on slide 245 ---------- #
    for shape in slide_245.shapes:
        click_action = shape.click_action  # always present (even if empty)

        # Determine the internal target slide (if any)
        target_slide = None
        try:
            target_slide = click_action.target_slide  # may raise if no rel exists
        except ValueError:
            target_slide = None

        target_idx = None
        if target_slide is not None:
            target_idx = slide_id_to_index.get(target_slide.slide_id)
            print(f"Found internal hyperlink on a shape to slide {target_idx + 1}")
        else:
            # External links are irrelevant for this task but still report for transparency
            if click_action.hyperlink and click_action.hyperlink.address:
                print(f"Found external hyperlink: {click_action.hyperlink.address}")

        # Retrieve visible text (if any) for later checks
        text_content = ""
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame and shape.text_frame.text:
            text_content = shape.text_frame.text.strip()

        # ----- 3a. Verify Logo link → slide 1 (index 0) ----- #
        if not logo_link_found and target_idx == 0:
            # Heuristic: logo usually doesn't contain the word "Next"
            if not re.search(r"next", text_content, re.IGNORECASE):
                logo_link_found = True
                print("✓ Logo correctly links to slide 1")

        # ----- 3b. Verify "Next" text box link → slide 2 (index 1) ----- #
        if not next_link_found and target_idx == 1:
            if re.search(r"\bnext\b", text_content, re.IGNORECASE):
                next_link_found = True
                print("✓ 'Next' text correctly links to slide 2")

    # ---------- 4. Scoring ---------- #
    if logo_link_found:
        score += 0.5
    else:
        print("✗ Logo hyperlink to slide 1 not found or incorrect")

    if next_link_found:
        score += 0.5
    else:
        print("✗ 'Next' hyperlink to slide 2 not found or incorrect")

    final_score = min(score, 1.0)
    print(f"REWARD: {final_score}")
    return final_score


# ----------------- Execute verification ----------------- #
if __name__ == "__main__":
    FILE_PATH = "/home/user/while_rehearsing_i_noticed_slide_245_doesnt_let_me_move_around_the_deck_smoothly_in_libreoffice_impr_golden.pptx"
    verify_hyperlinks(FILE_PATH)

