"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 229 of my LibreOffice Impress deck, Table 1 refuses to sit where I want it. Could you walk me through how to snap the table to the left side of the slide and apply a precise 1.0 cm left margin?
Generated: 2025-09-10 18:46:37
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
from pptx import Presentation

# ------------------------------------------------------------
# Reward Script: Verify Table Alignment on Slide 229
# ------------------------------------------------------------
# Task to verify
# "On slide 229 of my LibreOffice Impress deck, Table 1 refuses to sit where I want it. 
#  Walk me through how to snap the table to the left side of the slide and apply a precise 1.0 cm left margin."
# ------------------------------------------------------------
# Verification Logic
# 1. Load the provided PPTX file (no points – prerequisite)
# 2. Ensure slide 229 exists (no points – prerequisite)
# 3. Detect at least one table on slide 229        -> 0.5 points
# 4. Check each table’s left position equals 1.0 cm (360 000 EMU) 
#    within a small tolerance (20 000 EMU ≈ 0.22 cm) -> 0.5 points
# ------------------------------------------------------------
# Progressive scoring is used – partial credit awarded when only some
# requirements are fulfilled. Final score capped at 1.0.
# ------------------------------------------------------------

def verify_table_left_margin(file_path: str) -> float:
    """Verify that all tables on slide 229 are aligned to exactly 1 cm from the
    left edge of the slide.

    Parameters
    ----------
    file_path : str
        Absolute path to the presentation to verify.

    Returns
    -------
    float
        Reward between 0.0 and 1.0.
    """

    MAX_SCORE = 1.0
    score = 0.0

    # ---------------- Prerequisite Checks (NO POINTS) ----------------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Could not load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    slide_index = 228  # 0-based index for slide 229
    if slide_index >= len(prs.slides):
        print(f"✗ Slide 229 not found. Total slides: {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[slide_index]

    # ---------------- Requirement 1: Table presence ----------------
    table_shapes = [shp for shp in slide.shapes if shp.has_table]
    if not table_shapes:
        print("✗ No table found on slide 229 (0 points)")
    else:
        print(f"✓ Found {len(table_shapes)} table(s) on slide 229 (0.5 points)")
        score += 0.5

        # ---------------- Requirement 2: Correct left margin ----------------
        expected_left = 360_000  # 1 cm in EMU
        tolerance = 20_000       # ~0.22 cm tolerance
        all_tables_aligned = True

        for idx, tbl in enumerate(table_shapes, start=1):
            diff = abs(tbl.left - expected_left)
            print(f"  Table {idx}: left = {tbl.left} EMU, diff = {diff}")
            if diff <= tolerance:
                print("   ✓ Within tolerance")
            else:
                print("   ✗ Exceeds tolerance")
                all_tables_aligned = False

        if all_tables_aligned:
            print("✓ All tables aligned with a 1.0 cm left margin (0.5 points)")
            score += 0.5
        else:
            print("✗ Not all tables have the correct 1.0 cm left margin (0 points)")

    # ---------------- Final Score ----------------
    final_score = min(score, MAX_SCORE)
    print(f"REWARD: {final_score}")
    return final_score

# ------------------------------------------------------------
# Execute verification when run as a script
# ------------------------------------------------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/on_slide_229_of_my_libreoffice_impress_deck_table_1_refuses_to_sit_where_i_want_it_could_you_walk_me_golden.pptx"
    verify_table_left_margin(FILE_PATH)

