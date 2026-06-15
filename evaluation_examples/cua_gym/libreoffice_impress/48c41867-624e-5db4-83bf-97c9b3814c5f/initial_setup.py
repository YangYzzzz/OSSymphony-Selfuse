"""
Initial Setup: Create a 9-slide Benchmark Pitch presentation with slide 6 having title but empty content.
Task ID: impress_sales_077
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_077'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a text box with styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    # Standard 16:9 dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Benchmark Pitch"
    slide1.placeholders[1].text = "Q2 2025 Competitive Analysis\nPrepared by Strategic Sales Team"

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_textbox(slide2, Inches(0.8), Inches(0.5), Inches(11), Inches(1), "Executive Summary",
                font_size=32, bold=True, color=RGBColor(0x2B, 0x6C, 0xB0))
    add_textbox(slide2, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5),
                "Our enterprise platform has consistently outperformed market alternatives across "
                "all major benchmarks. In the past fiscal year, we achieved a 23% increase in customer "
                "satisfaction scores and reduced average deployment time by 40%. This presentation "
                "provides a comprehensive comparison of our solution against the two leading competitors, "
                "Competitor A (TechForward Solutions) and Competitor B (NexGen Systems), across five "
                "critical business dimensions.",
                font_size=16)

    # --- Slide 3: Our Product Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide3, Inches(0.8), Inches(0.5), Inches(11), Inches(1), "Our Product Overview",
                font_size=32, bold=True, color=RGBColor(0x2B, 0x6C, 0xB0))
    overview_text = (
        "Platform Highlights:\n\n"
        "- Cloud-native architecture with 99.97% uptime SLA\n"
        "- AI-powered analytics dashboard with real-time insights\n"
        "- Seamless integration with 150+ enterprise tools\n"
        "- Dedicated customer success managers for Enterprise tier\n"
        "- SOC 2 Type II and ISO 27001 certified\n"
        "- Average onboarding time: 14 business days"
    )
    add_textbox(slide3, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5), overview_text, font_size=16)

    # --- Slide 4: Market Analysis ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide4, Inches(0.8), Inches(0.5), Inches(11), Inches(1), "Market Analysis",
                font_size=32, bold=True, color=RGBColor(0x2B, 0x6C, 0xB0))
    market_text = (
        "The enterprise software market is projected to reach $650B by 2027, "
        "growing at a CAGR of 11.3%. Key trends shaping the competitive landscape:\n\n"
        "- Shift toward unified platforms over point solutions\n"
        "- Increasing demand for AI/ML-driven automation\n"
        "- Growing emphasis on security and compliance\n"
        "- Remote work driving cloud adoption acceleration\n\n"
        "Our market share grew from 12.4% to 15.8% in the past 12 months, "
        "while Competitor A held steady at 18.2% and Competitor B declined to 9.1%."
    )
    add_textbox(slide4, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5), market_text, font_size=16)

    # --- Slide 5: Customer Testimonials ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide5, Inches(0.8), Inches(0.5), Inches(11), Inches(1), "Customer Testimonials",
                font_size=32, bold=True, color=RGBColor(0x2B, 0x6C, 0xB0))
    testimonials = (
        '"After switching from Competitor A, our team productivity increased by 35%. '
        'The onboarding was seamless and support has been exceptional."\n'
        '   -- Sarah Mitchell, VP of Operations, Meridian Corp\n\n'
        '"We evaluated all three platforms extensively. The performance benchmarks '
        'spoke for themselves -- no other vendor came close on reliability."\n'
        '   -- David Park, CTO, Horizon Dynamics\n\n'
        '"The value proposition is unmatched. We reduced our total cost of ownership '
        'by 28% while gaining features our previous vendor couldn\'t offer."\n'
        '   -- Rachel Torres, Director of IT, Pinnacle Health Systems'
    )
    add_textbox(slide5, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5), testimonials, font_size=14)

    # --- Slide 6: "How We Stack Up" - EMPTY content (task target) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide6, Inches(0.8), Inches(0.5), Inches(11), Inches(1), "How We Stack Up",
                font_size=32, bold=True, color=RGBColor(0x2B, 0x6C, 0xB0))
    # No chart, no additional content - this is where the agent will add the chart

    # --- Slide 7: Implementation Timeline ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide7, Inches(0.8), Inches(0.5), Inches(11), Inches(1), "Implementation Timeline",
                font_size=32, bold=True, color=RGBColor(0x2B, 0x6C, 0xB0))
    timeline_text = (
        "Phase 1 - Discovery & Planning (Weeks 1-2)\n"
        "   Requirements gathering, stakeholder interviews, architecture review\n\n"
        "Phase 2 - Configuration & Integration (Weeks 3-5)\n"
        "   Platform setup, API integrations, data migration planning\n\n"
        "Phase 3 - Testing & Training (Weeks 6-7)\n"
        "   UAT, performance testing, end-user training sessions\n\n"
        "Phase 4 - Go-Live & Optimization (Weeks 8-10)\n"
        "   Phased rollout, monitoring, performance tuning, KPI tracking"
    )
    add_textbox(slide7, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5), timeline_text, font_size=16)

    # --- Slide 8: Pricing Plans ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide8, Inches(0.8), Inches(0.5), Inches(11), Inches(1), "Pricing Plans",
                font_size=32, bold=True, color=RGBColor(0x2B, 0x6C, 0xB0))
    # Add a simple table
    table_shape = slide8.shapes.add_table(4, 4, Inches(1), Inches(2), Inches(10), Inches(3))
    table = table_shape.table
    headers = ["Feature", "Starter", "Professional", "Enterprise"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    data = [
        ["Users", "Up to 25", "Up to 250", "Unlimited"],
        ["Storage", "50 GB", "500 GB", "5 TB"],
        ["Monthly Price", "$499/mo", "$1,999/mo", "Custom"],
    ]
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 9: Next Steps ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide9, Inches(0.8), Inches(0.5), Inches(11), Inches(1), "Next Steps",
                font_size=32, bold=True, color=RGBColor(0x2B, 0x6C, 0xB0))
    next_steps = (
        "1. Schedule a personalized demo with your technical team\n\n"
        "2. Request a 30-day proof-of-concept trial\n\n"
        "3. Review our detailed technical architecture documentation\n\n"
        "4. Connect with a customer reference in your industry\n\n"
        "Contact: sales@ourplatform.com | +1 (888) 555-0192\n"
        "Account Executive: James Whitfield | james.w@ourplatform.com"
    )
    add_textbox(slide9, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5), next_steps, font_size=16)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
