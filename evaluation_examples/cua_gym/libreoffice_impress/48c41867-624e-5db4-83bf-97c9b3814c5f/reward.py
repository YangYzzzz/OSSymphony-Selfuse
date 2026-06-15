"""
Reward Script: Insert horizontal grouped bar chart on slide 6 with competitive benchmarks
Task ID: impress_sales_077
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25) - Chart exists on slide 6 and is BAR_CLUSTERED (horizontal grouped bar)
  Component 2 (0.20) - Chart title is 'Competitive Benchmarks'
  Component 3 (0.20) - Correct 5 categories in order
  Component 4 (0.20) - Correct data values for all 3 series (Us, Competitor A, Competitor B)
  Component 5 (0.15) - Series 0 (Us) uses brand blue #2B6CB0
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_077'

# Expected data
EXPECTED_CATEGORIES = ['Performance', 'Reliability', 'Support', 'Value', 'Ease of Use']
EXPECTED_SERIES_VALUES = [
    [95.0, 99.0, 92.0, 88.0, 90.0],   # Us
    [82.0, 91.0, 75.0, 80.0, 78.0],   # Competitor A
    [70.0, 85.0, 68.0, 90.0, 82.0],   # Competitor B
]
BRAND_BLUE = '2B6CB0'


def persist_app_state():
    """Save any unsaved LibreOffice edits before verification."""
    try:
        os.environ["DISPLAY"] = ":0"
        import pyautogui
        import time
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[5]  # slide 6 (0-indexed)

    # Find chart shape on slide 6
    chart_shape = None
    for shape in slide.shapes:
        if shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Chart exists on slide 6 and is BAR_CLUSTERED (0.25 points)
    try:
        if chart_shape is None:
            print("FAIL: Component 1 -- No chart found on slide 6")
        else:
            chart = chart_shape.chart
            # BAR_CLUSTERED = 57 in python-pptx (horizontal grouped bar)
            chart_type_val = chart.chart_type
            if chart_type_val == 57:  # BAR_CLUSTERED
                print(f"PASS: Component 1 -- BAR_CLUSTERED chart found on slide 6 (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 -- Chart type is {chart_type_val}, expected BAR_CLUSTERED (57)")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Early exit if no chart found
    if chart_shape is None:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    chart = chart_shape.chart

    # Component 2: Chart title is 'Competitive Benchmarks' (0.20 points)
    try:
        if chart.has_title:
            actual_title = chart.chart_title.text_frame.text.strip()
            if actual_title == 'Competitive Benchmarks':
                print(f"PASS: Component 2 -- Chart title is 'Competitive Benchmarks' (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 2 -- Chart title is '{actual_title}', expected 'Competitive Benchmarks'")
        else:
            print("FAIL: Component 2 -- Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Correct 5 categories in order (0.20 points)
    try:
        plot = chart.plots[0]
        actual_cats = [str(c) for c in plot.categories]
        if actual_cats == EXPECTED_CATEGORIES:
            print(f"PASS: Component 3 -- All 5 categories match in order (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 3 -- Categories mismatch")
            print(f"  Expected: {EXPECTED_CATEGORIES}")
            print(f"  Actual:   {actual_cats}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Correct data values for all 3 series (0.20 points)
    try:
        num_series = len(chart.series)
        if num_series < 3:
            print(f"FAIL: Component 4 -- Only {num_series} series, expected 3")
        else:
            all_match = True
            for si in range(3):
                actual_vals = list(chart.series[si].values)
                expected_vals = EXPECTED_SERIES_VALUES[si]
                if actual_vals != expected_vals:
                    print(f"FAIL: Component 4 -- Series {si} values mismatch")
                    print(f"  Expected: {expected_vals}")
                    print(f"  Actual:   {actual_vals}")
                    all_match = False
            if all_match:
                print(f"PASS: Component 4 -- All 3 series have correct values (0.20 pts)")
                total_score += 0.20
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Series 0 (Us) uses brand blue #2B6CB0 (0.15 points)
    try:
        series0 = chart.series[0]
        fill = series0.format.fill
        if fill.type is not None:
            try:
                actual_rgb = str(fill.fore_color.rgb).upper()
                expected_rgb = BRAND_BLUE.upper()
                if actual_rgb == expected_rgb:
                    print(f"PASS: Component 5 -- Series 0 color is #{BRAND_BLUE} (0.15 pts)")
                    total_score += 0.15
                else:
                    print(f"FAIL: Component 5 -- Series 0 color is #{actual_rgb}, expected #{expected_rgb}")
            except Exception as e2:
                print(f"FAIL: Component 5 -- Could not read series 0 color: {e2}")
        else:
            print(f"FAIL: Component 5 -- Series 0 has no fill set (fill.type is None)")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
