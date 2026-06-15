"""
Reward Script: Insert three circular shapes with process steps and arrows on slide 3
Task ID: impress_sales_033
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25) - Three oval/circle shapes on slide 3
  Component 2 (0.25) - Correct text in ovals: '1 Discovery', '2 Implementation', '3 Success'
  Component 3 (0.25) - Brand blue fill (#2B6CB0) and white text on ovals
  Component 4 (0.25) - Arrow shapes connecting the circles
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_033'


def _safe_color_check(run):
    """Return True if run has a non-None RGB color type (safe to read .rgb)."""
    try:
        return run.font.color.type is not None
    except Exception:
        return False


def persist_app_state(domain):
    """Save any unsaved LibreOffice changes via Ctrl+S."""
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        import time
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[2]  # Slide 3 (0-indexed)

    # Collect oval shapes and arrow shapes on slide 3
    oval_shapes = []
    arrow_shapes = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                auto_type = shape.auto_shape_type
                # OVAL is type 9, but also accept ELLIPSE or similar circular shapes
                if auto_type is not None and auto_type == 9:  # MSO_AUTO_SHAPE_TYPE.OVAL
                    oval_shapes.append(shape)
                # RIGHT_ARROW is 33, LEFT_ARROW is 66, various arrow types
                # Accept any arrow-like auto shape
                elif auto_type is not None and 'ARROW' in str(auto_type):
                    arrow_shapes.append(shape)
            except Exception:
                pass
        # Also check for freeform or connector shapes that might serve as arrows
        elif shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            arrow_shapes.append(shape)

    print(f"INFO: Found {len(oval_shapes)} oval shapes and {len(arrow_shapes)} arrow shapes on slide 3")

    # Component 1: Three oval/circle shapes on slide 3 (0.25 points)
    try:
        if len(oval_shapes) >= 3:
            print(f"PASS: Component 1 — Found {len(oval_shapes)} oval shapes on slide 3 (0.25 pts)")
            total_score += 0.25
        elif len(oval_shapes) >= 1:
            partial = 0.25 * (len(oval_shapes) / 3.0)
            print(f"PARTIAL: Component 1 — Found {len(oval_shapes)}/3 oval shapes ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — No oval shapes found on slide 3")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Correct text labels in ovals (0.25 points)
    # Expected: '1 Discovery', '2 Implementation', '3 Success' (text may span paragraphs)
    expected_labels = {
        'discovery': False,
        'implementation': False,
        'success': False,
    }
    expected_numbers = {
        '1': False,
        '2': False,
        '3': False,
    }
    try:
        for shape in oval_shapes:
            if shape.has_text_frame:
                full_text = ' '.join(
                    para.text.strip() for para in shape.text_frame.paragraphs
                    if para.text.strip()
                ).lower()
                # Check for each label
                for label in expected_labels:
                    if label in full_text:
                        expected_labels[label] = True
                for num in expected_numbers:
                    if num in full_text:
                        expected_numbers[num] = True

        labels_found = sum(1 for v in expected_labels.values() if v)
        numbers_found = sum(1 for v in expected_numbers.values() if v)
        # Score: half for numbers, half for labels
        label_score = (labels_found / 3.0) * 0.125
        number_score = (numbers_found / 3.0) * 0.125
        comp2_score = label_score + number_score

        if comp2_score >= 0.25:
            print(f"PASS: Component 2 — All labels and numbers found ({comp2_score:.2f} pts)")
            total_score += comp2_score
        elif comp2_score > 0:
            missing_labels = [k for k, v in expected_labels.items() if not v]
            missing_nums = [k for k, v in expected_numbers.items() if not v]
            print(f"PARTIAL: Component 2 — Missing labels: {missing_labels}, missing numbers: {missing_nums} ({comp2_score:.2f} pts)")
            total_score += comp2_score
        else:
            print(f"FAIL: Component 2 — No correct text labels found in oval shapes")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Brand blue fill (#2B6CB0) and white text (0.25 points)
    try:
        blue_fill_count = 0
        white_text_count = 0
        for shape in oval_shapes:
            # Check fill color
            try:
                fill = shape.fill
                if fill.type is not None and fill.type == 1:  # SOLID fill
                    rgb = str(fill.fore_color.rgb).upper()
                    if rgb == '2B6CB0':
                        blue_fill_count += 1
            except Exception:
                pass

            # Check text color (white = FFFFFF)
            if shape.has_text_frame:
                white_found_in_shape = any(
                    str(run.font.color.rgb).upper() == 'FFFFFF'
                    for para in shape.text_frame.paragraphs
                    for run in para.runs
                    if run.text.strip() and _safe_color_check(run)
                )
                if white_found_in_shape:
                    white_text_count += 1

        # Need all 3 ovals to have blue fill and white text
        fill_score = min(blue_fill_count, 3) / 3.0 * 0.125
        text_color_score = min(white_text_count, 3) / 3.0 * 0.125
        comp3_score = fill_score + text_color_score

        if comp3_score >= 0.25:
            print(f"PASS: Component 3 — All ovals have #2B6CB0 fill and white text ({comp3_score:.2f} pts)")
            total_score += comp3_score
        elif comp3_score > 0:
            print(f"PARTIAL: Component 3 — {blue_fill_count}/3 blue fills, {white_text_count}/3 white text ({comp3_score:.2f} pts)")
            total_score += comp3_score
        else:
            print(f"FAIL: Component 3 — No ovals with correct fill/text colors")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Arrow shapes connecting the circles (0.25 points)
    try:
        if len(arrow_shapes) >= 2:
            print(f"PASS: Component 4 — Found {len(arrow_shapes)} arrow shapes on slide 3 (0.25 pts)")
            total_score += 0.25
        elif len(arrow_shapes) == 1:
            print(f"PARTIAL: Component 4 — Found 1/2 arrow shapes (0.125 pts)")
            total_score += 0.125
        else:
            print(f"FAIL: Component 4 — No arrow shapes found on slide 3")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
