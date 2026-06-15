"""
Initial Setup: Insert circular shapes on slide 3 with process steps
Task ID: impress_sales_033
Domain: libreoffice_impress

Creates a 7-slide sales presentation. Slide 3 has the title
'Our 3-Step Process' but an EMPTY content area (no shapes yet).
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_033'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q3 Sales Strategy"
    slide1.placeholders[1].text = "Accelerating Growth Through Innovation"

    # --- Slide 2: Market Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Market Overview"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Total addressable market grew 23% YoY to $4.7B"
    p = tf2.add_paragraph()
    p.text = "Key competitors: Apex Solutions, NovaTech, BridgePoint"
    p.level = 0
    p = tf2.add_paragraph()
    p.text = "Our market share increased from 8.2% to 11.5%"
    p.level = 0
    p = tf2.add_paragraph()
    p.text = "Enterprise segment shows strongest growth potential"
    p.level = 0

    # --- Slide 3: Our 3-Step Process (EMPTY content - task target) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    title_box = slide3.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(1.0)
    )
    tf3 = title_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Our 3-Step Process"
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.runs[0]
    run3.font.size = Pt(36)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    # Content area is intentionally empty - no circles, no arrows

    # --- Slide 4: Revenue Projections ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    title4 = slide4.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(1.0)
    )
    tf4t = title4.text_frame
    p4t = tf4t.paragraphs[0]
    p4t.text = "Revenue Projections"
    p4t.alignment = PP_ALIGN.CENTER
    r4t = p4t.runs[0]
    r4t.font.size = Pt(32)
    r4t.font.bold = True
    r4t.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    tbl_shape = slide4.shapes.add_table(
        5, 4, Inches(1.5), Inches(1.8), Inches(10), Inches(3.5)
    )
    tbl = tbl_shape.table
    headers = ["Region", "Q1 ($K)", "Q2 ($K)", "Q3 Target ($K)"]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    data = [
        ["North America", "1,245", "1,387", "1,520"],
        ["Europe", "876", "923", "1,050"],
        ["Asia Pacific", "654", "712", "890"],
        ["Latin America", "312", "345", "425"],
    ]
    for r, row in enumerate(data, 1):
        for c, val in enumerate(row):
            tbl.cell(r, c).text = val

    # --- Slide 5: Customer Testimonials ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    title5 = slide5.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(1.0)
    )
    tf5t = title5.text_frame
    p5t = tf5t.paragraphs[0]
    p5t.text = "What Our Customers Say"
    p5t.alignment = PP_ALIGN.CENTER
    r5t = p5t.runs[0]
    r5t.font.size = Pt(32)
    r5t.font.bold = True
    r5t.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    quotes = [
        ('"Their platform transformed our pipeline management. We saw a 40% increase in conversion rates within 3 months."',
         "- Sarah Chen, VP Sales at TechVentures"),
        ('"The onboarding process was seamless. Our team was fully productive in under two weeks."',
         "- Marcus Johnson, Director of Operations at Pinnacle Corp"),
    ]
    y_pos = Inches(1.8)
    for quote_text, attribution in quotes:
        qbox = slide5.shapes.add_textbox(
            Inches(1.5), y_pos, Inches(10), Inches(2.0)
        )
        qtf = qbox.text_frame
        qtf.word_wrap = True
        qp = qtf.paragraphs[0]
        qp.text = quote_text
        qp.alignment = PP_ALIGN.LEFT
        qrun = qp.runs[0]
        qrun.font.size = Pt(18)
        qrun.font.italic = True
        qrun.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

        ap = qtf.add_paragraph()
        ap.text = attribution
        ap.alignment = PP_ALIGN.RIGHT
        arun = ap.runs[0]
        arun.font.size = Pt(14)
        arun.font.bold = True
        arun.font.color.rgb = RGBColor(0x2B, 0x6C, 0xB0)

        y_pos += Inches(2.5)

    # --- Slide 6: Implementation Timeline ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    title6 = slide6.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(1.0)
    )
    tf6t = title6.text_frame
    p6t = tf6t.paragraphs[0]
    p6t.text = "Implementation Timeline"
    p6t.alignment = PP_ALIGN.CENTER
    r6t = p6t.runs[0]
    r6t.font.size = Pt(32)
    r6t.font.bold = True
    r6t.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    milestones = [
        ("Week 1-2", "Stakeholder alignment and requirements gathering"),
        ("Week 3-4", "Platform configuration and data migration"),
        ("Week 5-6", "Team training and pilot launch"),
        ("Week 7-8", "Full rollout and performance monitoring"),
    ]
    y = Inches(1.8)
    for period, desc in milestones:
        mbox = slide6.shapes.add_textbox(Inches(1.5), y, Inches(9), Inches(0.8))
        mtf = mbox.text_frame
        mtf.word_wrap = True
        mp = mtf.paragraphs[0]
        mr1 = mp.add_run()
        mr1.text = f"{period}:  "
        mr1.font.bold = True
        mr1.font.size = Pt(16)
        mr1.font.color.rgb = RGBColor(0x2B, 0x6C, 0xB0)
        mr2 = mp.add_run()
        mr2.text = desc
        mr2.font.size = Pt(16)
        mr2.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
        y += Inches(1.2)

    # --- Slide 7: Next Steps ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    title7 = slide7.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(1.0)
    )
    tf7t = title7.text_frame
    p7t = tf7t.paragraphs[0]
    p7t.text = "Next Steps"
    p7t.alignment = PP_ALIGN.CENTER
    r7t = p7t.runs[0]
    r7t.font.size = Pt(32)
    r7t.font.bold = True
    r7t.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    steps = [
        "Schedule follow-up meeting with leadership team by June 15",
        "Complete vendor evaluation scorecard for top 3 candidates",
        "Finalize budget allocation for Q3 technology investments",
        "Launch pilot program with North America enterprise accounts",
    ]
    sbox = slide7.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(4.0))
    stf = sbox.text_frame
    stf.word_wrap = True
    for i, step in enumerate(steps):
        if i == 0:
            sp = stf.paragraphs[0]
        else:
            sp = stf.add_paragraph()
        sp.text = step
        sp.level = 0
        srun = sp.runs[0]
        srun.font.size = Pt(18)
        srun.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
