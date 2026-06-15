"""
Initial Setup: Create Event.pptx presentation with 4 slides.
Slide 2 titled 'Store Launch' has NO Fontwork elements.
Task ID: impress_ndo_045
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_045'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard widescreen 10x7.5 inches
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Grand Opening Event"
    slide1.placeholders[1].text = "2025 Launch Series — Marketing Division"
    # Style title
    for run in slide1.shapes.title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # --- Slide 2: Store Launch (NO Fontwork) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Title text box
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Store Launch"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # Description text box
    txBox2 = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Join us for the grand opening of our newest flagship store in downtown Seattle."
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.runs[0]
    run2.font.name = "Arial"
    run2.font.size = Pt(20)
    run2.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)

    p3 = tf2.add_paragraph()
    p3.text = ""
    p4 = tf2.add_paragraph()
    p4.text = "Date: Saturday, June 14, 2025"
    p4.runs[0].font.name = "Arial"
    p4.runs[0].font.size = Pt(18)
    p4.runs[0].font.color.rgb = RGBColor(0x34, 0x49, 0x5E)
    p4.runs[0].font.bold = True

    p5 = tf2.add_paragraph()
    p5.text = "Time: 10:00 AM — 8:00 PM"
    p5.runs[0].font.name = "Arial"
    p5.runs[0].font.size = Pt(18)
    p5.runs[0].font.color.rgb = RGBColor(0x34, 0x49, 0x5E)

    p6 = tf2.add_paragraph()
    p6.text = "Location: 1425 3rd Avenue, Seattle, WA 98101"
    p6.runs[0].font.name = "Arial"
    p6.runs[0].font.size = Pt(18)
    p6.runs[0].font.color.rgb = RGBColor(0x34, 0x49, 0x5E)

    p7 = tf2.add_paragraph()
    p7.text = ""
    p8 = tf2.add_paragraph()
    p8.text = "Featuring live music, exclusive discounts, and complimentary refreshments for all attendees."
    p8.runs[0].font.name = "Arial"
    p8.runs[0].font.size = Pt(18)
    p8.runs[0].font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)

    # --- Slide 3: Schedule ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank

    txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
    tf3 = txBox3.text_frame
    p = tf3.paragraphs[0]
    p.text = "Event Schedule"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # Schedule table
    rows, cols = 6, 3
    tbl_shape = slide3.shapes.add_table(rows, cols, Inches(1.5), Inches(1.8), Inches(10), Inches(4))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(2.5)
    tbl.columns[1].width = Inches(4)
    tbl.columns[2].width = Inches(3.5)

    headers = ["Time", "Activity", "Location"]
    schedule_data = [
        ["10:00 AM", "Ribbon Cutting Ceremony", "Main Entrance"],
        ["11:00 AM", "Store Tour & Product Demos", "All Floors"],
        ["1:00 PM", "Live Jazz Performance", "Rooftop Terrace"],
        ["3:00 PM", "Customer Appreciation Raffle", "Ground Floor"],
        ["5:00 PM", "VIP Cocktail Reception", "Executive Lounge"],
    ]

    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    for r, row_data in enumerate(schedule_data, 1):
        for c, val in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(14)

    # --- Slide 4: Contact Information ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank

    txBox4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
    tf4 = txBox4.text_frame
    p = tf4.paragraphs[0]
    p.text = "Contact Information"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    contacts = [
        ("Event Coordinator:", "Priya Sharma — priya.sharma@company.com"),
        ("Marketing Lead:", "David Kim — david.kim@company.com"),
        ("Press Inquiries:", "media@company.com | (206) 555-0142"),
        ("RSVP:", "events.company.com/grand-opening"),
    ]

    txBox5 = slide4.shapes.add_textbox(Inches(2), Inches(2), Inches(9), Inches(4.5))
    tf5 = txBox5.text_frame
    tf5.word_wrap = True
    for i, (label, detail) in enumerate(contacts):
        if i > 0:
            tf5.add_paragraph().text = ""
        p_label = tf5.add_paragraph() if i > 0 or True else tf5.paragraphs[0]
        p_label.text = label
        p_label.runs[0].font.name = "Arial"
        p_label.runs[0].font.size = Pt(20)
        p_label.runs[0].font.bold = True
        p_label.runs[0].font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

        p_detail = tf5.add_paragraph()
        p_detail.text = detail
        p_detail.runs[0].font.name = "Arial"
        p_detail.runs[0].font.size = Pt(18)
        p_detail.runs[0].font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
