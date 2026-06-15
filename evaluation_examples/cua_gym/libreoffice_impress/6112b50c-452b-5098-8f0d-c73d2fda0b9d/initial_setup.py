"""
Initial Setup: Create a 10-slide presentation with no watermark on slide master
Task ID: impress_exec_082
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_082'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q1 2025 Strategic Review"
    slide1.placeholders[1].text = "Prepared by the Strategy & Operations Team"

    # Slide 2: Agenda
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Market Overview & Competitive Landscape"
    items2 = [
        "Revenue Performance by Region",
        "Product Development Pipeline Update",
        "Customer Acquisition Metrics",
        "Operational Efficiency Initiatives",
        "Financial Projections for Q2",
        "Risk Assessment & Mitigation",
        "Strategic Priorities & Action Items",
    ]
    for item in items2:
        p = body2.add_paragraph()
        p.text = item
        p.level = 0

    # Slide 3: Market Overview
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Market Overview"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Total addressable market grew 12.3% YoY to $84.7B"
    p = body3.add_paragraph()
    p.text = "Our market share increased from 8.2% to 9.1%"
    p = body3.add_paragraph()
    p.text = "Key competitor Nexion Corp lost 1.4 points to regulatory issues"
    p = body3.add_paragraph()
    p.text = "Emerging markets contributed 34% of total growth"

    # Slide 4: Revenue Performance
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Revenue Performance by Region"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "North America: $42.3M (+8.7%)"
    regions = [
        "Europe: $28.1M (+15.2%)",
        "Asia-Pacific: $19.6M (+22.4%)",
        "Latin America: $7.8M (+11.9%)",
        "Middle East & Africa: $3.4M (+6.3%)",
    ]
    for r in regions:
        p = body4.add_paragraph()
        p.text = r
        p.level = 0

    # Slide 5: Product Pipeline
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Product Development Pipeline"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Project Aurora - Enterprise analytics platform (Beta Q2)"
    items5 = [
        "Project Meridian - Mobile-first CRM integration (Alpha Q3)",
        "Project Zenith - AI-powered forecasting module (Design phase)",
        "Project Titan - Infrastructure modernization (70% complete)",
    ]
    for item in items5:
        p = body5.add_paragraph()
        p.text = item
        p.level = 0

    # Slide 6: Customer Metrics
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Customer Acquisition & Retention"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "New enterprise clients: 47 (target: 40)"
    items6 = [
        "Customer retention rate: 94.2% (up from 91.8%)",
        "Net Promoter Score: 72 (industry avg: 58)",
        "Average contract value: $185K (+12% YoY)",
        "Pipeline value: $312M across 234 opportunities",
    ]
    for item in items6:
        p = body6.add_paragraph()
        p.text = item
        p.level = 0

    # Slide 7: Operational Efficiency
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Operational Efficiency"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "Deployment cycle time reduced from 14 days to 6 days"
    items7 = [
        "Infrastructure costs down 23% through cloud optimization",
        "Support ticket resolution improved to 4.2 hours avg",
        "Automated testing coverage increased to 87%",
    ]
    for item in items7:
        p = body7.add_paragraph()
        p.text = item
        p.level = 0

    # Slide 8: Financial Projections
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Q2 2025 Financial Projections"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Projected revenue: $108.5M (+11.2% QoQ)"
    items8 = [
        "Gross margin target: 68.5%",
        "Operating expenses: $62.3M (within budget)",
        "EBITDA target: $31.7M",
        "Capital expenditure: $8.2M for data center expansion",
    ]
    for item in items8:
        p = body8.add_paragraph()
        p.text = item
        p.level = 0

    # Slide 9: Risk Assessment
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Risk Assessment"
    body9 = slide9.placeholders[1].text_frame
    body9.text = "Supply chain disruption risk: MEDIUM (mitigation in place)"
    items9 = [
        "Regulatory compliance risk: LOW (audit passed March 2025)",
        "Talent retention risk: MEDIUM (new compensation review Q2)",
        "Cybersecurity risk: LOW (SOC 2 Type II certified)",
        "Market volatility risk: HIGH (hedging strategy active)",
    ]
    for item in items9:
        p = body9.add_paragraph()
        p.text = item
        p.level = 0

    # Slide 10: Next Steps
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "Strategic Priorities & Next Steps"
    body10 = slide10.placeholders[1].text_frame
    body10.text = "Launch Project Aurora beta by June 15"
    items10 = [
        "Complete European expansion hiring (28 positions)",
        "Finalize partnership with Datastream Technologies",
        "Submit SOX compliance documentation by April 30",
        "Board strategy retreat: May 12-14, Lake Tahoe",
    ]
    for item in items10:
        p = body10.add_paragraph()
        p.text = item
        p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
