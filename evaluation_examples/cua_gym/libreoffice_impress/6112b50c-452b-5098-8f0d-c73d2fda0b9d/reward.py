"""
Reward Script: Set slide master watermark 'ACME CORP'
Task ID: impress_exec_082
Domain: libreoffice_impress
Scoring:
  Component 1: Watermark text 'ACME CORP' exists on slide master (0.25)
  Component 2: Font is 72pt Arial Bold (0.25)
  Component 3: Font color is #E8E8E8 (light gray) (0.20)
  Component 4: Rotation is 45 degrees (0.15)
  Component 5: Centered on slide (0.15)
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_082'


def find_watermark_shape(master):
    """Find a text shape on the slide master containing 'ACME CORP'."""
    for shape in master.shapes:
        if hasattr(shape, 'text') and 'ACME CORP' in (shape.text or '').upper():
            return shape
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Find the watermark shape on the slide master
    watermark = None
    for master in prs.slide_masters:
        watermark = find_watermark_shape(master)
        if watermark is not None:
            break

    if watermark is None:
        print("FAIL: No shape with 'ACME CORP' found on any slide master")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Found watermark shape: name='{watermark.name}', text='{watermark.text}'")

    # Component 1: Watermark text 'ACME CORP' exists on slide master (0.25 points)
    try:
        text_content = watermark.text.strip()
        if text_content == 'ACME CORP':
            print(f"PASS: Component 1 -- Watermark text is 'ACME CORP' (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 -- Expected 'ACME CORP', found '{text_content}'")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Font is 72pt Arial Bold (0.25 points)
    try:
        runs = []
        for para in watermark.text_frame.paragraphs:
            for run in para.runs:
                if (run.text or '').strip():
                    runs.append(run)

        if not runs:
            print("FAIL: Component 2 -- No text runs found in watermark shape")
        else:
            run = runs[0]
            font_name = run.font.name
            font_size = run.font.size  # in EMU
            font_bold = run.font.bold

            # 72pt = 72 * 12700 = 914400 EMU
            expected_size = 914400
            size_ok = (font_size is not None and abs(font_size - expected_size) < 12700)  # ~1pt tolerance
            name_ok = (font_name is not None and font_name.lower() == 'arial')
            bold_ok = (font_bold is True)

            comp2_score = 0.0
            if name_ok:
                comp2_score += 0.10
            if size_ok:
                comp2_score += 0.10
            if bold_ok:
                comp2_score += 0.05

            if name_ok and size_ok and bold_ok:
                print(f"PASS: Component 2 -- Font: {font_name} {font_size/12700:.0f}pt Bold={font_bold} (0.25 pts)")
                total_score += 0.25
            elif comp2_score > 0:
                print(f"PARTIAL: Component 2 -- Font: name={font_name}(ok={name_ok}), size={font_size}(ok={size_ok}), bold={font_bold}(ok={bold_ok}) ({comp2_score} pts)")
                total_score += comp2_score
            else:
                print(f"FAIL: Component 2 -- Font: name={font_name}(ok={name_ok}), size={font_size}(ok={size_ok}), bold={font_bold}(ok={bold_ok})")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Font color is #E8E8E8 (0.20 points)
    try:
        if runs:
            run = runs[0]
            try:
                color_type = run.font.color.type
                if color_type is not None:
                    color_rgb = str(run.font.color.rgb).upper()
                    if color_rgb == 'E8E8E8':
                        print(f"PASS: Component 3 -- Font color is #E8E8E8 (0.20 pts)")
                        total_score += 0.20
                    else:
                        print(f"FAIL: Component 3 -- Expected color #E8E8E8, found #{color_rgb}")
                else:
                    print(f"FAIL: Component 3 -- Font color type is None (no explicit color set)")
            except AttributeError:
                print(f"FAIL: Component 3 -- Could not read font color (theme-based or missing)")
        else:
            print("FAIL: Component 3 -- No runs to check color")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Rotation is 45 degrees (0.15 points)
    try:
        rotation = watermark.rotation
        if rotation is not None and abs(rotation - 45.0) < 1.0:
            print(f"PASS: Component 4 -- Rotation is {rotation} degrees (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 -- Expected rotation ~45 degrees, found {rotation}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Centered on slide (0.15 points)
    try:
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        shape_center_x = watermark.left + watermark.width // 2
        shape_center_y = watermark.top + watermark.height // 2
        slide_center_x = slide_w // 2
        slide_center_y = slide_h // 2

        # Tolerance: 10% of slide dimension
        tol_x = slide_w * 0.10
        tol_y = slide_h * 0.10

        centered_x = abs(shape_center_x - slide_center_x) <= tol_x
        centered_y = abs(shape_center_y - slide_center_y) <= tol_y

        if centered_x and centered_y:
            print(f"PASS: Component 5 -- Centered on slide (dx={abs(shape_center_x - slide_center_x)}, dy={abs(shape_center_y - slide_center_y)}) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 5 -- Not centered. shape_center=({shape_center_x},{shape_center_y}), slide_center=({slide_center_x},{slide_center_y})")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
