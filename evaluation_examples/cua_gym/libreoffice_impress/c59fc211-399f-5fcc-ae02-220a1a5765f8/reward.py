"""
FINAL REWARD SCRIPT - SUCCESS
Task: Insert page numbers in the footer, bottom-right, starting at page 1.
Generated: 2025-10-17 06:04:31
Status: success
Model: azure-o3
Total Steps: 9
"""

import os
from pptx import Presentation

# ---------------------------------------------
# Reward Verification Script
# Task: "Insert page numbers in the footer, bottom-right, starting at page 1."
# ---------------------------------------------
# Scoring rubric (progressive):
#   • 0.60 points – Every slide shows the correct sequential page number text (1, 2, 3 …)
#   • 0.40 points – That page-number text appears in the bottom-right footer area on every slide
#   → 1.0 only when BOTH conditions are satisfied for 100 % of slides
# ---------------------------------------------

# Helper ----------------------------------------------------------------------
def _is_bottom_right(shape, slide_w, slide_h):
    """Return True if the shape sits in the bottom-right quadrant (heuristic)."""
    left   = shape.left
    top    = shape.top
    right  = shape.left + shape.width
    bottom = shape.top  + shape.height

    # Treat bottom-right ~35 % of width/height as footer area
    horizontal_ok = left  >= slide_w * 0.65 or right  >= slide_w * 0.9
    vertical_ok   = top   >= slide_h * 0.65 or bottom >= slide_h * 0.9
    return horizontal_ok and vertical_ok

# Core verification -----------------------------------------------------------
def verify_presentation_page_numbers(file_path):
    """Return a reward score (0.0-1.0) verifying footer page numbers."""
    sequential_weight = 0.6
    position_weight   = 0.4

    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print("✗ Error opening presentation:", e)
        return 0.0

    total_slides = len(prs.slides)
    if total_slides == 0:
        print("✗ Presentation has no slides")
        return 0.0

    correct_numbers   = 0  # slides with correct numeric label
    correct_positions = 0  # slides where that label is bottom-right

    slide_w, slide_h = prs.slide_width, prs.slide_height

    for idx, slide in enumerate(prs.slides, start=1):
        number_found   = False
        position_found = False

        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.replace("\n", "").strip()

            # Accept pure digit or common variants like "Page 3" / "3 / N"
            if text == str(idx) or text.lower() == f"page {idx}" or text.startswith(f"{idx} /"):
                number_found = True
                if _is_bottom_right(shape, slide_w, slide_h):
                    position_found = True
                # No break: continue scanning in case multiple matches

        print(f"Slide {idx}: number_found={number_found}, position_ok={position_found}")

        if number_found:
            correct_numbers += 1
        if position_found:
            correct_positions += 1

    # Progressive scoring -----------------------------------------------------
    number_score   = sequential_weight * (correct_numbers   / total_slides)
    position_score = position_weight   * (correct_positions / total_slides)
    total_score    = round(min(number_score + position_score, 1.0), 3)

    print(f"✓ Correct sequential numbers   : {correct_numbers}/{total_slides} -> {number_score:.2f} points")
    print(f"✓ Correct bottom-right location: {correct_positions}/{total_slides} -> {position_score:.2f} points")
    print(f"Total reward score: {total_score}\n")

    return total_score

# -----------------------------------------------------------------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/insert_page_numbers_in_the_footer_bottom_right_starting_at_page_1.pptx"
    reward = verify_presentation_page_numbers(FILE_PATH)
    print(f"REWARD: {reward}")

