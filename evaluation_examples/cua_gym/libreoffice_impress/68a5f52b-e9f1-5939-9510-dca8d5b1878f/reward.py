"""
Reward Script: Insert a 4-column by 5-row table on slide 2 for student grades
Task ID: impress_tct_001
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Slide 2 contains a table shape
  Component 2 (0.3): Table dimensions are 4 columns x 5 rows
  Component 3 (0.3): Table has meaningful content (non-empty header and data cells)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_001'


def persist_app_state():
    """Save any unsaved LibreOffice edits before verification."""
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has only {len(prs.slides)} slide(s), need at least 2")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]  # 0-indexed, slide 2

    # Find table shapes on slide 2
    table_shapes = []
    for shape in slide2.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_shapes.append(shape)

    # Component 1: Slide 2 contains a table shape (0.4 points)
    try:
        if len(table_shapes) > 0:
            print(f"PASS: Component 1 -- Slide 2 contains {len(table_shapes)} table(s) (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 -- Slide 2 has no table shape")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    if len(table_shapes) == 0:
        # No table means nothing else to check
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    table = table_shapes[0].table

    # Component 2: Table dimensions are 4 columns x 5 rows (0.3 points)
    try:
        num_rows = len(table.rows)
        num_cols = len(table.columns)
        if num_rows == 5 and num_cols == 4:
            print(f"PASS: Component 2 -- Table is {num_cols} cols x {num_rows} rows (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 -- Expected 4 cols x 5 rows, found {num_cols} cols x {num_rows} rows")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Table has meaningful content -- non-empty header row and data cells (0.3 points)
    try:
        # Check that the header row (row 0) has non-empty cells
        header_cells = [table.cell(0, c).text.strip() for c in range(min(len(table.columns), 4))]
        non_empty_headers = [h for h in header_cells if h]

        # Check that at least some data rows have non-empty cells
        data_filled = 0
        for r in range(1, min(len(table.rows), 5)):
            row_texts = [table.cell(r, c).text.strip() for c in range(min(len(table.columns), 4))]
            non_empty = [t for t in row_texts if t]
            if len(non_empty) >= 2:
                data_filled += 1

        has_headers = len(non_empty_headers) >= 2
        has_data = data_filled >= 2

        if has_headers and has_data:
            print(f"PASS: Component 3 -- Table has {len(non_empty_headers)} headers and {data_filled} filled data rows (0.3 pts)")
            total_score += 0.3
        elif has_headers or has_data:
            partial = 0.15
            print(f"PARTIAL: Component 3 -- headers={has_headers}, data_rows={has_data} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 -- Table lacks meaningful content. Headers: {header_cells}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
