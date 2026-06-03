"""
Initial Setup: Insert a 4-column by 5-row table on slide 2
Task ID: impress_tct_001
Domain: libreoffice_impress

Creates a 3-slide presentation. Slide 2 has only a title 'Student Grades' and no table.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_001'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
    slide1.shapes.title.text = "Class Grades Report"
    slide1.placeholders[1].text = "Fall 2025 Semester\nPrepared by Ms. Rivera"

    # --- Slide 2: Title Only - "Student Grades" (NO table) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add a title text box at the top
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Student Grades"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x6E)

    # --- Slide 3: Summary slide ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Title
    txBox3_title = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf3_title = txBox3_title.text_frame
    p3_title = tf3_title.paragraphs[0]
    p3_title.text = "Summary & Next Steps"
    p3_title.alignment = PP_ALIGN.LEFT
    run3_title = p3_title.runs[0]
    run3_title.font.name = "Arial"
    run3_title.font.size = Pt(28)
    run3_title.font.bold = True
    run3_title.font.color.rgb = RGBColor(0x2E, 0x4A, 0x6E)

    # Body text
    txBox3_body = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
    tf3_body = txBox3_body.text_frame
    tf3_body.word_wrap = True
    p3a = tf3_body.paragraphs[0]
    p3a.text = "Overall class performance has been strong this semester."
    run3a = p3a.runs[0]
    run3a.font.name = "Arial"
    run3a.font.size = Pt(18)

    p3b = tf3_body.add_paragraph()
    p3b.text = "Parent-teacher conferences are scheduled for December 12-14."
    run3b = p3b.runs[0]
    run3b.font.name = "Arial"
    run3b.font.size = Pt(18)

    p3c = tf3_body.add_paragraph()
    p3c.text = "Final exams begin January 15, 2026."
    run3c = p3c.runs[0]
    run3c.font.name = "Arial"
    run3c.font.size = Pt(18)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
