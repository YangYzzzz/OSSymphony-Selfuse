"""
Initial Setup: Style Guide presentation with bullet list on slide 3
Task ID: impress_gf3_031
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_031'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_round_bullet(paragraph, char='●', color_hex=None):
    """Set bullet character on a paragraph."""
    pPr = paragraph._p.get_or_add_pPr()
    # Remove existing bullet elements
    for tag in [qn('a:buChar'), qn('a:buClr'), qn('a:buNone')]:
        for el in pPr.findall(tag):
            pPr.remove(el)
    # Add bullet char
    buChar = pPr.makeelement(qn('a:buChar'), {'char': char})
    pPr.append(buChar)
    if color_hex:
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgbClr = buClr.makeelement(qn('a:srgbClr'), {'val': color_hex})
        buClr.append(srgbClr)
        pPr.append(buClr)


def add_bullet_paragraph(text_frame, text, level=0, font_size=Pt(14),
                         bold=False, color=RGBColor(0, 0, 0), bullet_char='●'):
    """Add a bullet paragraph to the text frame."""
    p = text_frame.add_paragraph()
    p.level = level
    p.space_before = Pt(4)
    p.space_after = Pt(4)
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"
    set_round_bullet(p, char=bullet_char, color_hex='000000')
    return p


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Brand Style Guide"
    slide1.placeholders[1].text = "Nexus Digital Solutions — Q2 2025"

    # ---- Slide 2: Table of Contents ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Table of Contents"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    run.font.name = "Arial"

    toc_box = slide2.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    toc_tf = toc_box.text_frame
    toc_tf.word_wrap = True
    toc_items = [
        "1. Brand Voice & Messaging",
        "2. Typography Standards",
        "3. Content Formatting Guidelines",
        "4. Color Palette Usage",
        "5. Iconography & Visual Elements",
        "6. Document Templates",
        "7. Social Media Guidelines",
    ]
    for i, item in enumerate(toc_items):
        if i == 0:
            p = toc_tf.paragraphs[0]
        else:
            p = toc_tf.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        run.font.name = "Arial"
        p.space_before = Pt(8)

    # ---- Slide 3: Content Formatting Guidelines (THE KEY SLIDE) ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title
    title_box = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    ttf = title_box.text_frame
    p = ttf.paragraphs[0]
    p.text = "Content Formatting Guidelines"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    run.font.name = "Arial"

    # Bullet list text box
    bullet_box = slide3.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    btf = bullet_box.text_frame
    btf.word_wrap = True

    # First paragraph uses paragraphs[0] (already exists)
    p0 = btf.paragraphs[0]
    p0.level = 0
    p0.space_before = Pt(4)
    p0.space_after = Pt(4)
    run0 = p0.add_run()
    run0.text = "Use consistent heading hierarchy across all documents"
    run0.font.size = Pt(14)
    run0.font.bold = False
    run0.font.color.rgb = RGBColor(0, 0, 0)
    run0.font.name = "Arial"
    set_round_bullet(p0, char='●', color_hex='000000')

    # First-level item 2
    add_bullet_paragraph(btf,
        "Maintain minimum 1.15 line spacing for body text readability",
        level=0, font_size=Pt(14), bold=False,
        color=RGBColor(0, 0, 0), bullet_char='●')

    # Second-level sub-items under item 2
    add_bullet_paragraph(btf,
        "Exception: condensed layouts may use 1.0 spacing with approval",
        level=1, font_size=Pt(14), bold=False,
        color=RGBColor(0, 0, 0), bullet_char='●')

    add_bullet_paragraph(btf,
        "Footnotes and captions should use 0.9 spacing at smaller font sizes",
        level=1, font_size=Pt(14), bold=False,
        color=RGBColor(0, 0, 0), bullet_char='●')

    # First-level item 3
    add_bullet_paragraph(btf,
        "Apply brand color palette to charts, diagrams, and callout boxes",
        level=0, font_size=Pt(14), bold=False,
        color=RGBColor(0, 0, 0), bullet_char='●')

    # First-level item 4
    add_bullet_paragraph(btf,
        "Limit bullet point nesting to three levels maximum",
        level=0, font_size=Pt(14), bold=False,
        color=RGBColor(0, 0, 0), bullet_char='●')

    # ---- Slide 4: Color Palette Usage ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide4.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf4 = tb.text_frame
    p = tf4.paragraphs[0]
    p.text = "Color Palette Usage"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    run.font.name = "Arial"

    palette_box = slide4.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4))
    ptf = palette_box.text_frame
    ptf.word_wrap = True
    palette_text = [
        ("Primary Blue (#1A73E8)", "Used for headings, primary CTAs, and key highlights"),
        ("Secondary Teal (#00897B)", "Used for secondary UI elements and accents"),
        ("Neutral Dark (#333333)", "Standard body text and paragraph copy"),
        ("Neutral Light (#F5F5F5)", "Background fills and card surfaces"),
        ("Alert Red (#D32F2F)", "Error states, critical warnings, and urgent notices"),
    ]
    for i, (name, desc) in enumerate(palette_text):
        if i == 0:
            p = ptf.paragraphs[0]
        else:
            p = ptf.add_paragraph()
        run = p.add_run()
        run.text = f"{name}: {desc}"
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        run.font.name = "Arial"
        p.space_before = Pt(6)

    # ---- Slide 5: Iconography & Visual Elements ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    tb5 = slide5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf5 = tb5.text_frame
    p = tf5.paragraphs[0]
    p.text = "Iconography & Visual Elements"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    run.font.name = "Arial"

    body5 = slide5.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4.5))
    btf5 = body5.text_frame
    btf5.word_wrap = True
    p = btf5.paragraphs[0]
    run = p.add_run()
    run.text = ("Icons must follow the outlined style with 2px stroke weight. "
                "Filled icons are reserved for navigation elements and primary action buttons. "
                "Always maintain consistent sizing within icon groups — 24px for inline, "
                "48px for feature highlights, and 64px for hero sections.")
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    run.font.name = "Arial"

    # ---- Slide 6: Document Templates ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    tb6 = slide6.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf6 = tb6.text_frame
    p = tf6.paragraphs[0]
    p.text = "Document Templates"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    run.font.name = "Arial"

    body6 = slide6.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4.5))
    btf6 = body6.text_frame
    btf6.word_wrap = True
    templates = [
        "Quarterly Report — Use template QR-2025 with executive summary on page 1",
        "Client Proposal — Template CP-Standard includes pre-formatted pricing tables",
        "Internal Memo — Template IM-Brief for announcements under 500 words",
        "Meeting Notes — Template MN-Structured with action items section",
    ]
    for i, t in enumerate(templates):
        if i == 0:
            p = btf6.paragraphs[0]
        else:
            p = btf6.add_paragraph()
        run = p.add_run()
        run.text = t
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        run.font.name = "Arial"
        p.space_before = Pt(6)

    # ---- Slide 7: Social Media Guidelines ----
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    tb7 = slide7.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf7 = tb7.text_frame
    p = tf7.paragraphs[0]
    p.text = "Social Media Guidelines"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    run.font.name = "Arial"

    body7 = slide7.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4.5))
    btf7 = body7.text_frame
    btf7.word_wrap = True
    p = btf7.paragraphs[0]
    run = p.add_run()
    run.text = ("All social media posts must be reviewed by the brand communications team "
                "before publishing. Use the approved hashtag list and maintain a professional "
                "yet approachable tone. Image posts require the brand watermark in the "
                "bottom-right corner at 20% opacity.")
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    run.font.name = "Arial"

    # ---- Slide 8: Contact & Resources ----
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    tb8 = slide8.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf8 = tb8.text_frame
    p = tf8.paragraphs[0]
    p.text = "Contact & Resources"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    run.font.name = "Arial"

    body8 = slide8.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4.5))
    btf8 = body8.text_frame
    btf8.word_wrap = True
    contacts = [
        "Brand Team Lead: Priya Sharma — priya.sharma@nexusdigital.com",
        "Design Systems: Leo Martinez — leo.martinez@nexusdigital.com",
        "Asset Library: https://brand.nexusdigital.com/assets",
        "Template Repository: https://templates.nexusdigital.com",
        "Style Guide Feedback: styleguide@nexusdigital.com",
    ]
    for i, c in enumerate(contacts):
        if i == 0:
            p = btf8.paragraphs[0]
        else:
            p = btf8.add_paragraph()
        run = p.add_run()
        run.text = c
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        run.font.name = "Arial"
        p.space_before = Pt(6)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
