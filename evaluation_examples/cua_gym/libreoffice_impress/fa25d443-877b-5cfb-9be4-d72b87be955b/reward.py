"""
Reward Script: Verify bullet formatting changes on slide 3
Task ID: impress_gf3_031
Domain: libreoffice_impress
Scoring:
  C1 (0.20) - First-level bullet character is square (■)
  C2 (0.15) - First-level font size is 16pt (203200 EMU)
  C3 (0.15) - First-level text is bold
  C4 (0.20) - Second-level bullet character is em dash (—)
  C5 (0.10) - Second-level font size is 13pt (165100 EMU)
  C6 (0.20) - All bullet text color is #444444
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_031'


def persist_app_state(domain):
    """Save any unsaved LibreOffice changes before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(1.0)
            print("PERSIST: ctrl+s sent for libreoffice_impress")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify bullet formatting changes on slide 3.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[2]  # slide 3 (0-indexed)

    # Find the bullet text box (TextBox 3 or the shape with bullet paragraphs)
    bullet_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame:
            # Look for the shape that has bullet paragraphs (level 0 and level 1)
            paras = shape.text_frame.paragraphs
            has_bullets = False
            for p in paras:
                pPr = p._p.find('{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
                if pPr is not None:
                    bc = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
                    if bc is not None:
                        has_bullets = True
                        break
            if has_bullets:
                bullet_shape = shape
                break

    if bullet_shape is None:
        print("FAIL: No bullet text box found on slide 3")
        print("REWARD: 0.0")
        return 0.0

    # Separate paragraphs by level
    level0_paras = []
    level1_paras = []
    for para in bullet_shape.text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.level == 0:
            level0_paras.append(para)
        elif para.level == 1:
            level1_paras.append(para)

    print(f"Found {len(level0_paras)} first-level and {len(level1_paras)} second-level items")

    if len(level0_paras) == 0:
        print("FAIL: No first-level bullet items found")
        print("REWARD: 0.0")
        return 0.0

    # Helper to get buChar from paragraph XML
    def get_buchar(para):
        pPr = para._p.find('{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        if pPr is not None:
            bc = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
            if bc is not None:
                return bc.get('char')
        return None

    # Component 1: First-level bullet character is square (■) (0.20 points)
    try:
        all_square = True
        for para in level0_paras:
            bc = get_buchar(para)
            if bc != '■':
                all_square = False
                print(f"FAIL: C1 — level-0 para has buChar='{bc}', expected '■': {para.text[:40]}...")
                break
        if all_square:
            print(f"PASS: C1 — All {len(level0_paras)} first-level items have square bullet (■) (0.20 pts)")
            total_score += 0.20
    except Exception as e:
        print(f"ERROR: C1 — {e}")

    # Component 2: First-level font size is 16pt / 203200 EMU (0.15 points)
    try:
        all_correct_size = True
        for para in level0_paras:
            runs = [r for r in para.runs if (r.text or '').strip()]
            for run in runs:
                if run.font.size is None or run.font.size != 203200:
                    all_correct_size = False
                    actual = run.font.size
                    print(f"FAIL: C2 — level-0 run size={actual}, expected 203200 (16pt): {run.text[:30]}...")
                    break
            if not all_correct_size:
                break
        if all_correct_size:
            print(f"PASS: C2 — All first-level items have 16pt font size (0.15 pts)")
            total_score += 0.15
    except Exception as e:
        print(f"ERROR: C2 — {e}")

    # Component 3: First-level text is bold (0.15 points)
    try:
        all_bold = True
        for para in level0_paras:
            runs = [r for r in para.runs if (r.text or '').strip()]
            for run in runs:
                if run.font.bold is not True:
                    all_bold = False
                    print(f"FAIL: C3 — level-0 run bold={run.font.bold}, expected True: {run.text[:30]}...")
                    break
            if not all_bold:
                break
        if all_bold:
            print(f"PASS: C3 — All first-level items are bold (0.15 pts)")
            total_score += 0.15
    except Exception as e:
        print(f"ERROR: C3 — {e}")

    # Component 4: Second-level bullet character is em dash (—) (0.20 points)
    try:
        if len(level1_paras) == 0:
            print("FAIL: C4 — No second-level items found to check")
        else:
            all_emdash = True
            for para in level1_paras:
                bc = get_buchar(para)
                if bc != '\u2014':  # em dash
                    all_emdash = False
                    print(f"FAIL: C4 — level-1 para has buChar='{bc}', expected em dash: {para.text[:40]}...")
                    break
            if all_emdash:
                print(f"PASS: C4 — All {len(level1_paras)} second-level items have em dash bullet (0.20 pts)")
                total_score += 0.20
    except Exception as e:
        print(f"ERROR: C4 — {e}")

    # Component 5: Second-level font size is 13pt / 165100 EMU (0.10 points)
    try:
        if len(level1_paras) == 0:
            print("FAIL: C5 — No second-level items found")
        else:
            all_correct_size = True
            for para in level1_paras:
                runs = [r for r in para.runs if (r.text or '').strip()]
                for run in runs:
                    if run.font.size is None or run.font.size != 165100:
                        all_correct_size = False
                        actual = run.font.size
                        print(f"FAIL: C5 — level-1 run size={actual}, expected 165100 (13pt): {run.text[:30]}...")
                        break
                if not all_correct_size:
                    break
            if all_correct_size:
                print(f"PASS: C5 — All second-level items have 13pt font size (0.10 pts)")
                total_score += 0.10
    except Exception as e:
        print(f"ERROR: C5 — {e}")

    # Component 6: All bullet text color is #444444 (0.20 points)
    try:
        all_correct_color = True
        all_paras = level0_paras + level1_paras
        for para in all_paras:
            runs = [r for r in para.runs if (r.text or '').strip()]
            for run in runs:
                try:
                    if run.font.color.type is None:
                        all_correct_color = False
                        print(f"FAIL: C6 — run has no explicit color (inherited): {run.text[:30]}...")
                        break
                    rgb = str(run.font.color.rgb)
                    if rgb.upper() != '444444':
                        all_correct_color = False
                        print(f"FAIL: C6 — run color={rgb}, expected 444444: {run.text[:30]}...")
                        break
                except Exception:
                    all_correct_color = False
                    print(f"FAIL: C6 — cannot read color for run: {run.text[:30]}...")
                    break
            if not all_correct_color:
                break
        if all_correct_color:
            print(f"PASS: C6 — All bullet text has color #444444 (0.20 pts)")
            total_score += 0.20
    except Exception as e:
        print(f"ERROR: C6 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
