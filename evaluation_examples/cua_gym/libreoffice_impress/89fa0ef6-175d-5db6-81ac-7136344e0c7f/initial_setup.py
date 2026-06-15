"""
Initial Setup: Client Demo Presentation with Notes CSV
Task ID: osworld_multi_apps_impress_notes_import_014
Domain: libreoffice_impress

Creates:
  - /home/user/Client_Demo.pptx  (8 slides, empty notes panels)
  - /home/user/Desktop/demo_notes.csv  (slide_number, notes_text, time_allocation)
"""

import os
import csv
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
TASK_FILE = f'{WORKDIR}/Client_Demo.pptx'
CSV_FILE = f'{DESKTOP}/demo_notes.csv'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_presentation():
    """Create Client_Demo.pptx with 8 slides and NO notes."""
    prs = Presentation()

    slide_data = [
        {
            "layout": 0,
            "title": "Acme Corp Enterprise Solution",
            "subtitle": "Transforming Business Operations — Q2 2025",
        },
        {
            "layout": 1,
            "title": "Agenda",
            "body": (
                "• Company Overview\n"
                "• Current Challenges\n"
                "• Proposed Solution\n"
                "• Implementation Roadmap\n"
                "• Pricing & Packages\n"
                "• Case Studies\n"
                "• Next Steps"
            ),
        },
        {
            "layout": 1,
            "title": "Company Overview",
            "body": (
                "Founded in 2010 with a mission to simplify enterprise workflows\n\n"
                "• 500+ enterprise clients across 30 countries\n"
                "• $120M ARR with 35% YoY growth\n"
                "• ISO 27001 certified & SOC 2 Type II compliant\n"
                "• 24/7 dedicated support with < 1 hour SLA"
            ),
        },
        {
            "layout": 1,
            "title": "Current Challenges",
            "body": (
                "Key pain points we identified in your environment:\n\n"
                "• Manual data entry consuming ~18 hours/week per team\n"
                "• Siloed systems causing 30% reporting delays\n"
                "• Compliance audit prep taking 3 weeks per quarter\n"
                "• Onboarding new staff requires 6-week ramp-up"
            ),
        },
        {
            "layout": 1,
            "title": "Proposed Solution",
            "body": (
                "Acme Unified Platform — Core Features:\n\n"
                "• Automated data ingestion from 50+ source connectors\n"
                "• Real-time dashboards with role-based access control\n"
                "• One-click compliance reporting (GDPR, HIPAA, SOX)\n"
                "• AI-assisted onboarding with adaptive learning paths"
            ),
        },
        {
            "layout": 1,
            "title": "Implementation Roadmap",
            "body": (
                "Phase 1 (Weeks 1-4): Discovery & Configuration\n"
                "Phase 2 (Weeks 5-8): Data Migration & Integration\n"
                "Phase 3 (Weeks 9-12): UAT & Staff Training\n"
                "Phase 4 (Week 13+): Go-Live & Hypercare Support\n\n"
                "Dedicated project manager assigned throughout"
            ),
        },
        {
            "layout": 1,
            "title": "Pricing & Packages",
            "body": (
                "Starter — $2,500/month (up to 50 users)\n"
                "Professional — $6,000/month (up to 200 users)\n"
                "Enterprise — Custom pricing (unlimited users)\n\n"
                "All packages include:\n"
                "• Unlimited storage\n"
                "• 99.9% uptime SLA\n"
                "• Priority support"
            ),
        },
        {
            "layout": 1,
            "title": "Next Steps",
            "body": (
                "1. Sign LOI by April 30, 2025\n"
                "2. Schedule technical discovery call (May 5-9)\n"
                "3. Receive customised proposal within 5 business days\n"
                "4. Kick-off workshop — target May 19, 2025\n\n"
                "Contact: sarah.henderson@acmecorp.com | +1 (415) 555-0192"
            ),
        },
    ]

    for i, data in enumerate(slide_data):
        layout_idx = data.get("layout", 1)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = data["title"]

        # Set body / subtitle
        body_text = data.get("body") or data.get("subtitle", "")
        if body_text:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx != 0:  # not the title
                    ph.text = body_text
                    break

        # Ensure notes are EMPTY — do NOT set slide.notes_slide.notes_text_frame.text

    prs.save(TASK_FILE)
    print(f"Presentation created: {TASK_FILE}")


def create_csv():
    """Create demo_notes.csv on the Desktop."""
    os.makedirs(DESKTOP, exist_ok=True)

    notes_data = [
        (1, "Welcome attendees and introduce the presenting team. Briefly outline the session goals and expected outcomes.", 2),
        (2, "Walk through each agenda item. Confirm with the client if they want to adjust order or add discussion points.", 3),
        (3, "Highlight our global footprint and compliance credentials. Emphasise the 500+ client milestone and growth trajectory.", 5),
        (4, "Reference the discovery workshop findings. Use specific metrics gathered from their IT and ops teams.", 2),
        (5, "Demo the live dashboard if projector allows. Focus on the compliance reporting module as the client flagged this as critical.", 3),
        (6, "Stress that all four phases have fixed-price milestones. Offer to share the detailed GANTT chart post-meeting.", 4),
        (7, "Be prepared for negotiation on Professional tier. Discount authority up to 10% without escalation.", 2),
        (8, "Close with energy — confirm decision timeline. Leave business cards and schedule follow-up before leaving the room.", 4),
    ]

    with open(CSV_FILE, 'w', newline='', encoding='utf-8') as f:
        writer = csv.writer(f)
        writer.writerow(['slide_number', 'notes_text', 'time_allocation'])
        for row in notes_data:
            writer.writerow(row)

    print(f"CSV file created: {CSV_FILE}")


def main():
    create_presentation()
    create_csv()

    # GUI-ready startup: open Client_Demo.pptx in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{TASK_FILE}"', delay_sec=3.0)
    print("GUI_READY: launched LibreOffice Impress with DISPLAY=:0")


main()
