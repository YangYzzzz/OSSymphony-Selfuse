"""
Reward Script: Import CSV notes into LibreOffice Impress slides
Task ID: osworld_multi_apps_impress_notes_import_014
Domain: libreoffice_impress
Scoring:
  Component 1: All 8 slides have non-empty notes text       (0.3 pts)
  Component 2: Each note contains the '(Time: Xmin)' suffix (0.3 pts)
  Component 3: Full note content matches CSV data exactly    (0.4 pts)
Total: 1.0
"""

import os

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_multi_apps_impress_notes_import_014'

# Expected data derived from the task: slide_number, notes_text, time_allocation
# Source: demo_notes.csv on Desktop
EXPECTED_NOTES = [
    (1, "Welcome attendees and introduce the presenting team. Briefly outline the session goals and expected outcomes.", 2),
    (2, "Walk through each agenda item. Confirm with the client if they want to adjust order or add discussion points.", 3),
    (3, "Highlight our global footprint and compliance credentials. Emphasise the 500+ client milestone and growth trajectory.", 5),
    (4, "Reference the discovery workshop findings. Use specific metrics gathered from their IT and ops teams.", 2),
    (5, "Demo the live dashboard if projector allows. Focus on the compliance reporting module as the client flagged this as critical.", 3),
    (6, "Stress that all four phases have fixed-price milestones. Offer to share the detailed GANTT chart post-meeting.", 4),
    (7, "Be prepared for negotiation on Professional tier. Discount authority up to 10% without escalation.", 2),
    (8, "Close with energy -- confirm decision timeline. Leave business cards and schedule follow-up before leaving the room.", 4),
]


def get_slide_notes(slide):
    """Extract notes text from a slide safely."""
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def verify_task(file_path):
    """
    Verify task completion: CSV notes inserted into presentation slides with time suffix.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have 8 slides
    if len(prs.slides) != 8:
        print(f"CRITICAL: Expected 8 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Read actual notes from all slides
    actual_notes = []
    for i, slide in enumerate(prs.slides):
        notes_text = get_slide_notes(slide)
        actual_notes.append(notes_text)
        print(f"Slide {i+1} notes: {repr(notes_text[:80])}")

    # Component 1: All 8 slides have non-empty notes (0.3 points)
    # This FAILS on initial (all empty) and PASSES on golden (all populated)
    try:
        non_empty_count = sum(1 for n in actual_notes if n != "")
        if non_empty_count == 8:
            print(f"PASS: Component 1 — All 8 slides have non-empty notes (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 — Only {non_empty_count}/8 slides have non-empty notes")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Each note contains the '(Time: Xmin)' suffix (0.3 points)
    # This FAILS on initial (all empty) and PASSES on golden (all have time suffix)
    try:
        time_suffix_count = 0
        for i, (slide_num, notes_text, time_alloc) in enumerate(EXPECTED_NOTES):
            expected_suffix = f"(Time: {time_alloc}min)"
            if expected_suffix in actual_notes[i]:
                time_suffix_count += 1
            else:
                print(f"FAIL: Component 2 — Slide {slide_num} missing '{expected_suffix}', got: {repr(actual_notes[i][-30:])}")
        if time_suffix_count == 8:
            print(f"PASS: Component 2 — All 8 slides have correct '(Time: Xmin)' suffix (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 — Only {time_suffix_count}/8 slides have correct time suffix")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Full note content matches expected (notes_text + ' (Time: Xmin)') exactly (0.4 points)
    # This FAILS on initial (empty notes) and PASSES on golden (correct full content)
    try:
        exact_match_count = 0
        for i, (slide_num, notes_text, time_alloc) in enumerate(EXPECTED_NOTES):
            expected_full = f"{notes_text} (Time: {time_alloc}min)"
            actual = actual_notes[i]
            if actual == expected_full:
                exact_match_count += 1
            else:
                print(f"FAIL: Component 3 — Slide {slide_num} mismatch:")
                print(f"  Expected: {repr(expected_full)}")
                print(f"  Actual:   {repr(actual)}")
        if exact_match_count == 8:
            print(f"PASS: Component 3 — All 8 slides match exact expected content (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 3 — Only {exact_match_count}/8 slides have exact content match")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/Client_Demo.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
