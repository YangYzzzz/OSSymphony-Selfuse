"""
Initial Setup: Create Annual_Review_2025.pptx with 12 slides.
Slide 9 has title '2026 Strategic Roadmap' and is otherwise empty.
Task ID: impress_ps_019
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_019'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.text = body_lines[0]
    for line in body_lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
    return slide


def add_blank_slide_with_title(prs, title_text):
    """Add a slide using Title Only layout (index 5) with just a title."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    # Add title as a text box at the top
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    add_title_slide(prs, "Annual Review 2025", "Meridian Technologies Inc.\nBoard of Directors Presentation")

    # Slide 2: Agenda
    add_content_slide(prs, "Agenda", [
        "Financial Performance Overview",
        "Product & Engineering Milestones",
        "Customer Growth & Retention Metrics",
        "Talent & Organizational Highlights",
        "Market Landscape Analysis",
        "Risk Assessment & Compliance Update",
        "2026 Strategic Roadmap",
        "Capital Allocation & Budget Proposal",
        "Key Decisions for Board Approval",
    ])

    # Slide 3: Financial Performance
    add_content_slide(prs, "Financial Performance Overview", [
        "Annual revenue reached $187.3M, up 34% year-over-year",
        "Gross margin improved to 72.1% from 68.4% in 2024",
        "Operating expenses held at 54% of revenue (down from 61%)",
        "Free cash flow positive for first time: $12.8M",
        "ARR crossed $200M run-rate in Q4 2025",
    ])

    # Slide 4: Product & Engineering
    add_content_slide(prs, "Product & Engineering Milestones", [
        "Shipped v4.0 platform with real-time collaboration engine",
        "Launched analytics dashboard with 15+ custom report templates",
        "Reduced average API latency from 340ms to 89ms",
        "99.97% uptime achieved across all production services",
        "Filed 8 patents in ML-driven workflow automation",
    ])

    # Slide 5: Customer Growth
    add_content_slide(prs, "Customer Growth & Retention", [
        "Total customers grew to 2,847 (up from 1,920)",
        "Net revenue retention rate: 118%",
        "Enterprise segment (>$100K ARR): 194 accounts, +62% YoY",
        "NPS score improved to 67 (industry avg: 41)",
        "Churn reduced to 4.2% annually (from 7.1%)",
    ])

    # Slide 6: Talent & Organization
    add_content_slide(prs, "Talent & Organizational Highlights", [
        "Headcount grew to 412 employees across 6 offices",
        "Engineering team expanded to 178 (43% of company)",
        "Employee satisfaction score: 4.3/5.0 (Gallup Q12)",
        "Voluntary attrition dropped to 8.7% (industry avg: 15%)",
        "Opened Singapore office to support APAC expansion",
    ])

    # Slide 7: Market Landscape
    add_content_slide(prs, "Market Landscape Analysis", [
        "Total addressable market estimated at $14.2B by 2027",
        "Three new competitors entered mid-market segment in 2025",
        "Regulatory tailwinds: EU Digital Services Act driving adoption",
        "AI copilot integrations becoming table stakes in our category",
        "Partnership pipeline: 12 channel partners signed in H2 2025",
    ])

    # Slide 8: Risk & Compliance
    add_content_slide(prs, "Risk Assessment & Compliance Update", [
        "SOC 2 Type II certification renewed (zero findings)",
        "GDPR data processing audit completed successfully",
        "Cyber insurance coverage increased to $25M",
        "Supply chain risk: AWS dependency mitigated with multi-cloud strategy",
        "Key person risk: succession plans documented for C-suite",
    ])

    # Slide 9: 2026 Strategic Roadmap — EMPTY except title
    slide9 = add_blank_slide_with_title(prs, "2026 Strategic Roadmap")

    # Slide 10: Capital Allocation
    add_content_slide(prs, "Capital Allocation & Budget Proposal", [
        "Proposed 2026 operating budget: $142M (+18% over 2025)",
        "R&D investment: $58M (41% of budget) for platform and AI",
        "Sales & marketing: $47M to support enterprise push",
        "G&A efficiency target: reduce to 12% of revenue",
        "Capital reserve: maintain $30M minimum cash position",
    ])

    # Slide 11: Board Decisions
    add_content_slide(prs, "Key Decisions for Board Approval", [
        "Approve 2026 operating budget of $142M",
        "Authorize Series D bridge financing ($40M convertible note)",
        "Greenlight APAC expansion with Singapore hub investment",
        "Approve executive equity refresh pool (2.5% of diluted shares)",
        "Ratify updated data governance and AI ethics policy",
    ])

    # Slide 12: Thank You / Q&A
    add_title_slide(prs, "Thank You", "Questions & Discussion\nconfidential@meridiantech.com")

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
