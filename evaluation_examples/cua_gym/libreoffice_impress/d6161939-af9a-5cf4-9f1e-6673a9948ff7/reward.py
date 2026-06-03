"""
Reward Script: Strategy roadmap slide with horizontal timeline and initiative blocks
Task ID: impress_ps_019
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.20): Four quarter labels Q1-Q4 present on slide 9
  - Component 2 (0.50): All 6 initiative blocks with correct names
  - Component 3 (0.15): Initiative blocks are colored shapes (solid fill)
  - Component 4 (0.15): Initiative block text font size ~12pt
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_019'

# The 6 required initiatives (match by substring to be resilient to minor text variations)
REQUIRED_INITIATIVES = [
    'Platform Redesign',
    'Mobile App Launch',
    'APAC Expansion',
    'AI Features',
    'Enterprise Tier',
    'IPO Prep',
]

REQUIRED_QUARTERS = ['Q1', 'Q2', 'Q3', 'Q4']


def get_all_text_shapes(slide):
    """Recursively get all shapes with text, including group children."""
    def extract(shape):
        results = []
        if hasattr(shape, 'text') and hasattr(shape, 'text_frame'):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.extend(extract(sub))
        return results
    out = []
    for shape in slide.shapes:
        out.extend(extract(shape))
    return out


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Check slide count (must have at least 9 slides)
    if len(prs.slides) < 9:
        print(f"CRITICAL: Presentation has only {len(prs.slides)} slides, need at least 9")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[8]  # 0-indexed, slide 9
    all_shapes = get_all_text_shapes(slide)

    # Collect all text content and shapes for analysis
    text_shapes = []
    for shape in all_shapes:
        text = shape.text.strip() if shape.text else ''
        if text:
            text_shapes.append((shape, text))

    # Component 1: Quarter labels Q1-Q4 present (0.20 points)
    # These must be separate text elements (not part of initiative names)
    try:
        found_quarters = set()
        for shape, text in text_shapes:
            clean = text.strip()
            if clean in REQUIRED_QUARTERS:
                found_quarters.add(clean)

        quarter_count = len(found_quarters)
        if quarter_count == 4:
            print(f"PASS: Component 1 — All 4 quarter labels found: {sorted(found_quarters)} (0.20 pts)")
            total_score += 0.20
        elif quarter_count > 0:
            partial = round(0.20 * quarter_count / 4, 2)
            print(f"PARTIAL: Component 1 — {quarter_count}/4 quarter labels found: {sorted(found_quarters)} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — No quarter labels (Q1-Q4) found on slide 9")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: All 6 initiative blocks present with correct names (0.50 points)
    # Each initiative is worth ~0.083 points
    try:
        found_initiatives = []
        for initiative in REQUIRED_INITIATIVES:
            if any(initiative.lower() in text.lower() for _, text in text_shapes):
                found_initiatives.append(initiative)

        init_count = len(found_initiatives)
        if init_count == 6:
            print(f"PASS: Component 2 — All 6 initiatives found (0.50 pts)")
            total_score += 0.50
        elif init_count > 0:
            partial = round(0.50 * init_count / 6, 2)
            print(f"PARTIAL: Component 2 — {init_count}/6 initiatives found: {found_initiatives} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No initiative blocks found on slide 9")

        missing = [i for i in REQUIRED_INITIATIVES if i not in found_initiatives]
        if missing:
            print(f"  Missing initiatives: {missing}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Initiative blocks are colored shapes with solid fill (0.15 points)
    # Check that initiative shapes are AUTO_SHAPE (rectangles) with solid fill color
    try:
        colored_initiative_count = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                shape_text = shape.text.strip() if hasattr(shape, 'text') and shape.text else ''
                is_initiative = any(init.lower() in shape_text.lower() for init in REQUIRED_INITIATIVES)
                if is_initiative:
                    try:
                        fill = shape.fill
                        if fill.type is not None and fill.type == 1:  # SOLID fill
                            colored_initiative_count += 1
                    except Exception:
                        pass
            # Also check group shapes
            if hasattr(shape, 'shapes'):
                for sub in shape.shapes:
                    if sub.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        sub_text = sub.text.strip() if hasattr(sub, 'text') and sub.text else ''
                        is_initiative = any(init.lower() in sub_text.lower() for init in REQUIRED_INITIATIVES)
                        if is_initiative:
                            try:
                                fill = sub.fill
                                if fill.type is not None and fill.type == 1:
                                    colored_initiative_count += 1
                            except Exception:
                                pass

        if colored_initiative_count >= 6:
            print(f"PASS: Component 3 — {colored_initiative_count} initiative blocks have colored fills (0.15 pts)")
            total_score += 0.15
        elif colored_initiative_count > 0:
            partial = round(0.15 * colored_initiative_count / 6, 2)
            print(f"PARTIAL: Component 3 — {colored_initiative_count}/6 initiative blocks with colored fills ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No initiative blocks with colored fills found")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Initiative text font size approximately 12pt (152400 EMU) (0.15 points)
    # Task says "12pt font" for initiative block text
    try:
        correct_font_count = 0
        for shape, text in text_shapes:
            is_initiative = any(init.lower() in text.lower() for init in REQUIRED_INITIATIVES)
            if is_initiative:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size is not None:
                            # 12pt = 152400 EMU. Allow tolerance of +/- 2pt (25400 EMU per pt)
                            if abs(run.font.size - 152400) <= 50800:
                                correct_font_count += 1
                                break
                    else:
                        continue
                    break

        if correct_font_count >= 6:
            print(f"PASS: Component 4 — {correct_font_count} initiatives have ~12pt font (0.15 pts)")
            total_score += 0.15
        elif correct_font_count > 0:
            partial = round(0.15 * correct_font_count / 6, 2)
            print(f"PARTIAL: Component 4 — {correct_font_count}/6 initiatives with ~12pt font ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — No initiative blocks with ~12pt font found")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
