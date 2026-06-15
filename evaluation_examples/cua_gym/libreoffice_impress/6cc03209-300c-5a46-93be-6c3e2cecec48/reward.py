"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 5 is still using the default dark font, and it’s getting lost against my navy backdrop. In LibreOffice Impress, how can I quickly change every text box on that single slide to the bright yellow color #FFFF00?
Generated: 2025-09-10 13:10:04
Status: success
Model: azure-o3
Total Steps: 2
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
import os


def verify_slide5_yellow(file_path: str) -> float:
    """Verify that every non-empty text run on slide 5 is bright yellow (#FFFF00).

    Scoring (progressive):
        1.0 – 100% of text runs are bright yellow
        0.8 – ≥ 80% of text runs are bright yellow
        0.5 – ≥ 50% of text runs are bright yellow
        0.2 – ≥ 20% of text runs are bright yellow
        0.0 – < 20% or verification failed
    """

    print(f"Verifying presentation file: {file_path}")

    # ------------------------------------------------------------------
    # Preliminary checks (no points awarded for these natural conditions)
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File not found – cannot verify task.")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        return 0.0

    if len(prs.slides) < 5:
        print(f"✗ Expected at least 5 slides, found {len(prs.slides)} – slide 5 missing.")
        return 0.0

    # ------------------------------------------------------------------
    # Core verification – inspect colours of all text runs on slide 5
    # ------------------------------------------------------------------
    slide = prs.slides[4]
    print("✓ Slide 5 found. Inspecting text colours…")

    bright_yellow = RGBColor(255, 255, 0)
    total_runs = 0
    yellow_runs = 0

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue  # Only interested in text-bearing shapes

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                # Ignore purely empty runs that sometimes appear in PPTX files
                if not run.text or not run.text.strip():
                    continue

                total_runs += 1
                rgb = run.font.color.rgb  # May be None if colour not explicitly set

                print(f"    Run text: '{run.text.strip()[:40]}' → Colour: {rgb}")

                if rgb == bright_yellow:
                    yellow_runs += 1

    # ------------------------------------------------------------------
    # Scoring logic – award points based on proportion of yellow text runs
    # ------------------------------------------------------------------
    if total_runs == 0:
        print("✗ No text runs found on slide 5 – nothing to verify.")
        return 0.0

    proportion = yellow_runs / total_runs
    print(f"Total non-empty text runs: {total_runs}")
    print(f"Bright yellow runs (#FFFF00): {yellow_runs}")
    print(f"Yellow proportion: {proportion:.2%}")

    if proportion == 1.0:
        score = 1.0
    elif proportion >= 0.8:
        score = 0.8
    elif proportion >= 0.5:
        score = 0.5
    elif proportion >= 0.2:
        score = 0.2
    else:
        score = 0.0

    print(f"Calculated score: {score}")
    return score


if __name__ == "__main__":
    test_file = "/home/user/slide_5_is_still_using_the_default_dark_font_and_its_getting_lost_against_my_navy_backdrop_in_libreo_golden.pptx"
    reward = verify_slide5_yellow(test_file)
    print(f"REWARD: {reward}")

