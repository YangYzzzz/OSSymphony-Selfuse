"""
Initial Setup: Create Action_Items presentation with 4 slides, slide 2 has 6 bullet points.
Task ID: impress_tct_075
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_075'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Action Items"
    slide1.placeholders[1].text = "Q2 2025 Sprint Planning"

    # --- Slide 2: Bullet list with 6 items (default round bullets) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    txBox = slide2.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8.0), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    # Title line
    p_title = tf.paragraphs[0]
    p_title.text = "Sprint Action Items"
    p_title.alignment = PP_ALIGN.LEFT
    run_title = p_title.runs[0]
    run_title.font.size = Pt(28)
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    bullet_items = [
        "Review and finalize the database migration plan with the backend team",
        "Update the client dashboard to include real-time analytics widgets",
        "Schedule user acceptance testing sessions for the new payment module",
        "Prepare the quarterly security audit report for stakeholder review",
        "Coordinate with the design team on the mobile app redesign mockups",
        "Set up automated CI/CD pipeline for the microservices architecture",
    ]

    from pptx.oxml.ns import qn

    for item_text in bullet_items:
        p = tf.add_paragraph()
        p.text = item_text
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        # Add default round bullet character
        pPr = p._p.get_or_add_pPr()
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
        pPr.append(buChar)

    # --- Slide 3: Timeline slide ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8.0), Inches(5.5))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True

    p3_title = tf3.paragraphs[0]
    p3_title.text = "Project Timeline"
    run3t = p3_title.runs[0]
    run3t.font.size = Pt(28)
    run3t.font.bold = True
    run3t.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    milestones = [
        ("April 7", "Database migration kick-off"),
        ("April 14", "Dashboard prototype ready"),
        ("April 21", "UAT sessions begin"),
        ("April 28", "Security audit submission"),
    ]
    for date, desc in milestones:
        p = tf3.add_paragraph()
        p.text = f"{date} - {desc}"
        p.level = 0
        run = p.runs[0]
        run.font.size = Pt(16)

    # --- Slide 4: Team Responsibilities ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8.0), Inches(5.5))
    tf4 = txBox4.text_frame
    tf4.word_wrap = True

    p4_title = tf4.paragraphs[0]
    p4_title.text = "Team Responsibilities"
    run4t = p4_title.runs[0]
    run4t.font.size = Pt(28)
    run4t.font.bold = True
    run4t.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    members = [
        ("Elena Rodriguez", "Backend Lead - Database migration oversight"),
        ("James Park", "Frontend Lead - Dashboard and analytics"),
        ("Aisha Patel", "QA Lead - UAT coordination and test plans"),
        ("David Chen", "DevOps Lead - CI/CD pipeline setup"),
        ("Sophie Larsson", "Design Lead - Mobile app redesign"),
    ]
    for name, role in members:
        p = tf4.add_paragraph()
        r1 = p.add_run()
        r1.text = name
        r1.font.size = Pt(16)
        r1.font.bold = True
        r1.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
        r2 = p.add_run()
        r2.text = f" - {role}"
        r2.font.size = Pt(16)
        r2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
