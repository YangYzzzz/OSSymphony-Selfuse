"""
Reward Script: Convert bullet list on slide 2 to numbered list (1), 2), etc.) with bold numbers
Task ID: impress_tct_075
Domain: libreoffice_impress
Scoring:
  Component 1: Numbered prefix "N) " present on all 6 items (0.4 pts)
  Component 2: Number prefix runs are bold (0.3 pts)
  Component 3: Original bullet characters removed (0.15 pts)
  Component 4: Original text content preserved (0.15 pts)
"""

import os
import re
import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_075'

# Expected original text content for the 6 list items (without any prefix)
EXPECTED_TEXTS = [
    "Review and finalize the database migration plan with the backend team",
    "Update the client dashboard to include real-time analytics widgets",
    "Schedule user acceptance testing sessions for the new payment module",
    "Prepare the quarterly security audit report for stakeholder review",
    "Coordinate with the design team on the mobile app redesign mockups",
    "Set up automated CI/CD pipeline for the microservices architecture",
]

NUM_ITEMS = 6


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"CRITICAL: Expected at least 2 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[1]  # slide 2 (0-indexed)

    # Find the text box with the list (not the title placeholder)
    list_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts = [p.text for p in shape.text_frame.paragraphs]
            # The list shape has multiple paragraphs with content
            non_empty = [t for t in texts if t.strip()]
            if len(non_empty) >= 6:
                list_shape = shape
                break

    if list_shape is None:
        print("CRITICAL: Could not find text box with 6+ items on slide 2")
        print("REWARD: 0.0")
        return 0.0

    # Get the list paragraphs (skip the title paragraph "Sprint Action Items")
    all_paras = list_shape.text_frame.paragraphs
    list_paras = []
    for p in all_paras:
        text = p.text.strip()
        if text and text != "Sprint Action Items":
            list_paras.append(p)

    if len(list_paras) < NUM_ITEMS:
        print(f"CRITICAL: Expected {NUM_ITEMS} list items, found {len(list_paras)}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Numbered prefix "N) " present on all 6 items (0.4 points)
    try:
        prefix_count = 0
        for idx, para in enumerate(list_paras[:NUM_ITEMS]):
            expected_prefix = f"{idx + 1})"
            full_text = para.text.strip()
            if full_text.startswith(expected_prefix):
                prefix_count += 1
                print(f"  PASS: Item {idx+1} starts with '{expected_prefix}'")
            else:
                print(f"  FAIL: Item {idx+1} expected prefix '{expected_prefix}', got: '{full_text[:15]}'")

        if prefix_count == NUM_ITEMS:
            print(f"PASS: Component 1 - All {NUM_ITEMS} items have correct numbered prefix (0.4 pts)")
            total_score += 0.4
        elif prefix_count > 0:
            partial = round(0.4 * prefix_count / NUM_ITEMS, 2)
            print(f"PARTIAL: Component 1 - {prefix_count}/{NUM_ITEMS} items have correct prefix ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 - No items have numbered prefix")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Number prefix runs are bold (0.3 points)
    try:
        bold_count = 0
        for idx, para in enumerate(list_paras[:NUM_ITEMS]):
            expected_prefix = f"{idx + 1}) "
            runs = para.runs
            if len(runs) == 0:
                print(f"  FAIL: Item {idx+1} has no runs")
                continue

            # The first run should contain the number prefix and be bold
            first_run = runs[0]
            first_text = first_run.text
            is_bold = first_run.font.bold is True

            if expected_prefix in first_text and is_bold:
                bold_count += 1
                print(f"  PASS: Item {idx+1} prefix run '{first_text}' is bold")
            elif expected_prefix.strip() in first_text and is_bold:
                # Allow minor spacing differences
                bold_count += 1
                print(f"  PASS: Item {idx+1} prefix run '{first_text}' is bold (minor spacing diff)")
            else:
                # Check if any run contains the prefix and is bold
                matched_any = any(
                    re.match(r'^\d+\)\s*', run.text) and run.font.bold is True
                    for run in runs
                )
                if matched_any:
                    bold_count += 1
                    print(f"  PASS: Item {idx+1} has bold number prefix in alternate run")
                else:
                    print(f"  FAIL: Item {idx+1} first run text='{first_text}', bold={first_run.font.bold}")

        if bold_count == NUM_ITEMS:
            print(f"PASS: Component 2 - All {NUM_ITEMS} number prefixes are bold (0.3 pts)")
            total_score += 0.3
        elif bold_count > 0:
            partial = round(0.3 * bold_count / NUM_ITEMS, 2)
            print(f"PARTIAL: Component 2 - {bold_count}/{NUM_ITEMS} prefixes are bold ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 - No number prefixes are bold")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Original bullet characters removed (0.15 points)
    try:
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        bullets_remaining = 0

        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide2.xml') as f:
                root = ET.parse(f).getroot()
                paras_xml = root.findall('.//a:p', ns)
                for para_xml in paras_xml:
                    text = ''.join(t.text or '' for t in para_xml.findall('.//a:t', ns))
                    if not text.strip() or text.strip() == "Sprint Action Items":
                        continue
                    pPr = para_xml.find('a:pPr', ns)
                    if pPr is not None:
                        buChar = pPr.find('a:buChar', ns)
                        if buChar is not None:
                            bullets_remaining += 1

        if bullets_remaining == 0:
            print(f"PASS: Component 3 - No bullet characters remain on list items (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 3 - {bullets_remaining} bullet characters still present")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Numbered prefix present AND original text content preserved (0.15 points)
    # This is a compound check: the numbered prefix must be present (task-introduced change)
    # AND the original text after the prefix must match. This ensures the component only
    # passes on golden (where numbers were added) and fails on initial (no numbers).
    try:
        preserved_count = 0
        for idx, para in enumerate(list_paras[:NUM_ITEMS]):
            full_text = para.text.strip()
            expected_prefix = f"{idx + 1})"
            # Must start with the numbered prefix (task change) AND have correct content
            if full_text.startswith(expected_prefix):
                content = re.sub(r'^\d+\)\s*', '', full_text)
                expected = EXPECTED_TEXTS[idx]
                if content == expected:
                    preserved_count += 1
                else:
                    print(f"  FAIL: Item {idx+1} content mismatch: got '{content[:50]}' expected '{expected[:50]}'")
            else:
                print(f"  FAIL: Item {idx+1} missing numbered prefix, cannot verify content preservation")

        if preserved_count == NUM_ITEMS:
            print(f"PASS: Component 4 - All {NUM_ITEMS} items have numbered prefix with preserved text (0.15 pts)")
            total_score += 0.15
        elif preserved_count > 0:
            partial = round(0.15 * preserved_count / NUM_ITEMS, 2)
            print(f"PARTIAL: Component 4 - {preserved_count}/{NUM_ITEMS} items correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 - No items have numbered prefix with preserved text")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    # Try with _initial suffix
    alt_path = f'{WORKDIR}/{TASK_ID}_initial.pptx'
    if os.path.exists(alt_path):
        file_path = alt_path
    else:
        print(f"File not found: {file_path}")
        print("REWARD: 0.0")
        import sys
        sys.exit(0)

verify_task(file_path)
