"""
Reward Script: Reverse the order of slides 3-8 in History.pptx
Task ID: impress_ndo_074
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.60): Slides 3-8 titles are in reversed order (Era 6..Era 1), ~0.10 per correct slide
  - Component 2 (0.25): Body content of reversed slides is preserved (moved with titles)
  - Component 3 (0.15): Boundary slides (1-2, 9-10) remain unchanged
  Precondition gate: File loads, has 10 slides
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_074'

# Expected titles for slides 3-8 AFTER reversal
EXPECTED_REVERSED_TITLES = ['Era 6', 'Era 5', 'Era 4', 'Era 3', 'Era 2', 'Era 1']

# Expected boundary slide titles (slides 1,2,9,10 - 0-indexed: 0,1,8,9)
BOUNDARY_TITLES = {
    0: 'History of Computing',
    1: 'Table of Contents',
    8: 'Timeline Summary',
    9: 'References',
}

# Key body text snippets for each Era slide (used to verify content moved with title)
# These are distinctive first bullet items from each Era slide
ERA_BODY_SNIPPETS = {
    'Era 1': 'Blaise Pascal invented the Pascaline',
    'Era 2': 'Colossus at Bletchley Park cracked German Lorenz ciphers',
    'Era 3': 'Bell Labs invented the transistor',
    'Era 4': 'Intel 4004 was the first commercial microprocessor',
    'Era 5': 'IBM PC standardized the personal computer platform',
    'Era 6': 'Google Search launched in 1998',
}


def get_slide_title(slide):
    """Extract the first non-empty text from a slide (title)."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    return text
    return ''


def get_slide_body_text(slide):
    """Get all text from a slide as a single string."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    texts.append(text)
    return ' '.join(texts)


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition gate: must have exactly 10 slides
    num_slides = len(prs.slides)
    if num_slides != 10:
        print(f"PRECONDITION FAIL: Expected 10 slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)

    # Component 1: Slides 3-8 titles are in reversed order (0.60 points, 0.10 per correct title)
    try:
        correct_titles = 0
        for i in range(6):  # slides 3-8 are indices 2-7
            slide_idx = i + 2
            actual_title = get_slide_title(slides[slide_idx])
            expected_title = EXPECTED_REVERSED_TITLES[i]
            if actual_title == expected_title:
                print(f"PASS: Slide {slide_idx + 1} title is '{actual_title}' (expected '{expected_title}')")
                correct_titles += 1
            else:
                print(f"FAIL: Slide {slide_idx + 1} title is '{actual_title}', expected '{expected_title}'")

        title_score = correct_titles * 0.10
        if title_score > 0:
            total_score += title_score
        print(f"Component 1 subtotal: {title_score:.2f}/0.60 ({correct_titles}/6 titles correct)")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Body content preserved after reversal (0.25 points)
    # Check that the body text of each reversed slide contains the expected snippet
    try:
        correct_bodies = 0
        for i in range(6):
            slide_idx = i + 2
            actual_title = get_slide_title(slides[slide_idx])
            body_text = get_slide_body_text(slides[slide_idx])

            # The expected snippet corresponds to the title on this slide
            if actual_title in ERA_BODY_SNIPPETS:
                expected_snippet = ERA_BODY_SNIPPETS[actual_title]
                if expected_snippet in body_text:
                    print(f"PASS: Slide {slide_idx + 1} body contains expected content for '{actual_title}'")
                    correct_bodies += 1
                else:
                    print(f"FAIL: Slide {slide_idx + 1} body missing expected snippet for '{actual_title}'")
            else:
                print(f"FAIL: Slide {slide_idx + 1} has unexpected title '{actual_title}', cannot verify body")

        # Only award body points if at least some titles were correct (otherwise meaningless)
        if correct_titles > 0:
            body_score = (correct_bodies / 6) * 0.25
            total_score += body_score
            print(f"Component 2 subtotal: {body_score:.2f}/0.25 ({correct_bodies}/6 bodies correct)")
        else:
            print(f"Component 2: Skipped (no titles matched, body check not meaningful)")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Boundary slides unchanged (0.15 points)
    # This component is scored only when reversal has happened (at least 1 title changed from initial order)
    # Initial order would be Era 1..Era 6, reversed is Era 6..Era 1
    # If slide 3 title is NOT 'Era 1' (initial), reversal has been attempted
    try:
        slide3_title = get_slide_title(slides[2])
        reversal_attempted = (slide3_title != 'Era 1')

        if reversal_attempted:
            correct_boundaries = 0
            for idx, expected_title in BOUNDARY_TITLES.items():
                actual_title = get_slide_title(slides[idx])
                if actual_title == expected_title:
                    correct_boundaries += 1
                    print(f"PASS: Boundary slide {idx + 1} title is '{actual_title}'")
                else:
                    print(f"FAIL: Boundary slide {idx + 1} title is '{actual_title}', expected '{expected_title}'")

            boundary_score = (correct_boundaries / 4) * 0.15
            if boundary_score > 0:
                total_score += boundary_score
            print(f"Component 3 subtotal: {boundary_score:.2f}/0.15 ({correct_boundaries}/4 boundaries correct)")
        else:
            print(f"Component 3: No reversal detected, boundary check yields 0 (initial state)")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
