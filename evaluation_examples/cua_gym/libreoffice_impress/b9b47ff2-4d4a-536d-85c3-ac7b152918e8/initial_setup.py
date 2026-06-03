"""
Initial Setup: Reverse the order of slides 3 through 8
Task ID: impress_ndo_074
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_074'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    """Add a title slide (layout 0)."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    """Add a content slide (layout 1) with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def create_initial():
    prs = Presentation()

    # Slide 1: Title Slide
    add_title_slide(
        prs,
        "History of Computing",
        "From Mechanical Calculators to Quantum Processors"
    )

    # Slide 2: Table of Contents
    add_content_slide(prs, "Table of Contents", [
        "Era 1: Mechanical Computation (1600s-1930s)",
        "Era 2: Vacuum Tubes and Early Mainframes (1940s-1950s)",
        "Era 3: Transistors and Minicomputers (1960s)",
        "Era 4: Integrated Circuits and Microprocessors (1970s)",
        "Era 5: Personal Computing Revolution (1980s-1990s)",
        "Era 6: Internet and Mobile Age (2000s-Present)",
    ])

    # Slide 3: Era 1
    add_content_slide(prs, "Era 1", [
        "Blaise Pascal invented the Pascaline in 1642",
        "Gottfried Leibniz built the Step Reckoner in 1694",
        "Charles Babbage designed the Difference Engine in 1822",
        "Ada Lovelace wrote the first algorithm for the Analytical Engine",
        "Herman Hollerith's tabulating machine processed the 1890 US Census",
    ])

    # Slide 4: Era 2
    add_content_slide(prs, "Era 2", [
        "Colossus at Bletchley Park cracked German Lorenz ciphers (1943)",
        "ENIAC weighed 30 tons and used 18,000 vacuum tubes (1945)",
        "UNIVAC I predicted the 1952 presidential election results",
        "Grace Hopper developed the first compiler (A-0) in 1952",
        "IBM 701 became the first commercially successful scientific computer",
    ])

    # Slide 5: Era 3
    add_content_slide(prs, "Era 3", [
        "Bell Labs invented the transistor in 1947, commercialized in the 1960s",
        "DEC PDP-8 was the first successful minicomputer ($18,500 in 1965)",
        "IBM System/360 introduced compatible computer family architecture",
        "Douglas Engelbart demonstrated the mouse and hypertext in 1968",
        "ARPANET connected UCLA and Stanford Research Institute in 1969",
    ])

    # Slide 6: Era 4
    add_content_slide(prs, "Era 4", [
        "Intel 4004 was the first commercial microprocessor (1971)",
        "Altair 8800 launched the hobbyist computing movement (1975)",
        "Apple II brought color graphics to home computing (1977)",
        "Xerox PARC pioneered the graphical user interface concept",
        "Ethernet networking was developed by Robert Metcalfe at PARC",
    ])

    # Slide 7: Era 5
    add_content_slide(prs, "Era 5", [
        "IBM PC standardized the personal computer platform (1981)",
        "Apple Macintosh popularized the GUI for consumers (1984)",
        "Microsoft Windows 3.1 reached 10 million sales (1992)",
        "Linus Torvalds released the Linux kernel in 1991",
        "Tim Berners-Lee invented the World Wide Web at CERN (1989)",
    ])

    # Slide 8: Era 6
    add_content_slide(prs, "Era 6", [
        "Google Search launched in 1998, transforming information access",
        "Apple iPhone revolutionized mobile computing in 2007",
        "Cloud computing enabled on-demand infrastructure scaling",
        "Deep learning breakthroughs led to modern AI assistants",
        "Quantum computing reached milestones with Google Sycamore (2019)",
    ])

    # Slide 9: Timeline Summary
    add_content_slide(prs, "Timeline Summary", [
        "1642-1930s: Mechanical computation foundations",
        "1940s-1950s: Electronic computing with vacuum tubes",
        "1960s: Transistor-based minicomputers emerge",
        "1970s: Microprocessor and integrated circuit revolution",
        "1980s-1990s: Personal computing goes mainstream",
        "2000s-Present: Internet, mobile, AI, and quantum era",
    ])

    # Slide 10: References
    add_content_slide(prs, "References", [
        'Ceruzzi, P.E. "A History of Modern Computing" (MIT Press, 2003)',
        'Campbell-Kelly, M. "Computer: A History of the Information Machine" (2014)',
        'Isaacson, W. "The Innovators" (Simon & Schuster, 2014)',
        'Turing, A. "On Computable Numbers" (1936)',
        'IEEE Annals of the History of Computing, various issues',
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
