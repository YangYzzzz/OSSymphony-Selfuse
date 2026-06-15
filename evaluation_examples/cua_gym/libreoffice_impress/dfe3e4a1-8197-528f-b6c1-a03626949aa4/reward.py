"""
Reward Script: Apply italic formatting to all subtitle text across the presentation.
Task ID: impstruct_021
Domain: libreoffice_impress
Scoring:
  Component 1 (0.6): Proportion of subtitle placeholders with italic runs (0.1 per slide, 6 slides)
  Component 2 (0.4): All subtitles italic AND no titles accidentally made italic (compound check)

Strategy: Identify title vs subtitle by placeholder idx (0 vs 1). If only idx=0 placeholders
exist (LibreOffice resave), fall back to shape order on each slide (first=title, second=subtitle).
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impstruct_021'

# Known subtitle texts (for fallback identification)
KNOWN_SUBTITLES = [
    "Annual Research Summary - Q4 2025 Findings",
    "Mixed-methods approach combining surveys and behavioral tracking",
    "Analysis of 12,400 respondents across 8 metropolitan regions",
    "Price sensitivity decreased by 14% compared to previous fiscal year",
    "Platform engagement correlates with conversion rates at r=0.73",
    "Strategic priorities for the upcoming product launch cycle",
]

KNOWN_TITLES = [
    "Consumer Behavior in Digital Markets",
    "Methodology Overview",
    "Key Demographic Insights",
    "Purchase Decision Factors",
    "Social Media Influence on Buying",
    "Conclusions and Recommendations",
]


def persist_app_state(domain):
    """Save any unsaved LibreOffice edits before verifying."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_shape_text(shape):
    """Get the full text of a shape."""
    if not shape.has_text_frame:
        return ""
    return "".join(
        run.text or ""
        for para in shape.text_frame.paragraphs
        for run in para.runs
    ).strip()


def classify_shapes(slide):
    """
    Classify placeholder shapes into title and subtitle.
    Returns (title_shape, subtitle_shape) or (None, None).
    Strategy: Use ph_idx if distinct (0 vs 1). Otherwise, match by known text.
    """
    placeholders = []
    for shape in slide.shapes:
        if shape.is_placeholder and shape.has_text_frame:
            placeholders.append(shape)

    if len(placeholders) < 2:
        return None, None

    # Check if we have distinct ph_idx values
    idx_map = {}
    for ph in placeholders:
        idx = ph.placeholder_format.idx
        if idx not in idx_map:
            idx_map[idx] = []
        idx_map[idx].append(ph)

    if 0 in idx_map and 1 in idx_map:
        # Standard structure
        return idx_map[0][0], idx_map[1][0]

    # Fallback: match by known text content
    title_shape = None
    subtitle_shape = None
    for ph in placeholders:
        txt = get_shape_text(ph)
        if txt in KNOWN_TITLES:
            title_shape = ph
        elif txt in KNOWN_SUBTITLES:
            subtitle_shape = ph

    if title_shape and subtitle_shape:
        return title_shape, subtitle_shape

    # Last fallback: first = title, second = subtitle by order
    if len(placeholders) >= 2:
        return placeholders[0], placeholders[1]

    return None, None


def check_all_runs_italic(shape):
    """Check if all non-empty runs in the shape are italic."""
    runs = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if (run.text or "").strip():
                runs.append(run)
    if not runs:
        return False, []
    italic_values = [r.font.italic for r in runs]
    all_italic = all(v is True for v in italic_values)
    return all_italic, italic_values


def check_no_runs_italic(shape):
    """Check that no non-empty runs in the shape are italic (i.e., italic is None or False)."""
    runs = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if (run.text or "").strip():
                runs.append(run)
    if not runs:
        return True, []
    italic_values = [r.font.italic for r in runs]
    none_italic = all(v is not True for v in italic_values)
    return none_italic, italic_values


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)
    if len(slides) == 0:
        print("CRITICAL: Presentation has no slides")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Proportion of subtitle placeholders with all runs italic (0.6 points)
    # Each slide contributes 0.1 points (6 slides max = 0.6)
    try:
        italic_subtitle_count = 0
        subtitle_count = 0
        per_slide_points = 0.6 / max(len(slides), 1)

        for i, slide in enumerate(slides):
            title_shape, subtitle_shape = classify_shapes(slide)
            if subtitle_shape is None:
                print(f"  Slide {i+1}: No subtitle found, skipping")
                continue
            subtitle_count += 1
            all_italic, italic_vals = check_all_runs_italic(subtitle_shape)
            if all_italic:
                italic_subtitle_count += 1
                total_score += per_slide_points
                print(f"  PASS: Slide {i+1} subtitle is italic ({per_slide_points:.2f} pts)")
            else:
                print(f"  FAIL: Slide {i+1} subtitle italic values: {italic_vals}")

        print(f"Component 1: {italic_subtitle_count}/{subtitle_count} subtitles italic "
              f"({min(total_score, 0.6):.2f}/0.60 pts)")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: All subtitles italic AND no titles accidentally made italic (0.4 points)
    # This compound check requires the task change (all subtitles italic) as prerequisite,
    # so it FAILS on initial_env where subtitles are not italic.
    try:
        all_subtitles_italic = (italic_subtitle_count == subtitle_count and subtitle_count > 0)
        italic_title_count = 0  # count titles that were incorrectly italicized

        for i, slide in enumerate(slides):
            title_shape, _ = classify_shapes(slide)
            if title_shape is None:
                continue
            no_italic, italic_vals = check_no_runs_italic(title_shape)
            if not no_italic:
                italic_title_count += 1
                print(f"  WARN: Slide {i+1} title has italic runs (should not): {italic_vals}")

        titles_untouched = (italic_title_count == 0)  # derived from actual check

        if all_subtitles_italic and titles_untouched:
            print(f"PASS: Component 2 - All subtitles italic, no titles affected (0.40 pts)")
            total_score += 0.4
        elif all_subtitles_italic and not titles_untouched:
            print(f"PARTIAL: Component 2 - Subtitles italic but some titles also italic (0.20 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 2 - Not all subtitles are italic "
                  f"(subtitles_italic={all_subtitles_italic}, titles_untouched={titles_untouched})")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
