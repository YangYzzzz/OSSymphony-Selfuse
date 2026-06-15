"""
Initial Setup: Create a 6-slide research findings presentation with non-italic subtitles
Task ID: impstruct_021
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impstruct_021'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide data: (title, subtitle)
    slides_data = [
        (
            "Consumer Behavior in Digital Markets",
            "Annual Research Summary - Q4 2025 Findings"
        ),
        (
            "Methodology Overview",
            "Mixed-methods approach combining surveys and behavioral tracking"
        ),
        (
            "Key Demographic Insights",
            "Analysis of 12,400 respondents across 8 metropolitan regions"
        ),
        (
            "Purchase Decision Factors",
            "Price sensitivity decreased by 14% compared to previous fiscal year"
        ),
        (
            "Social Media Influence on Buying",
            "Platform engagement correlates with conversion rates at r=0.73"
        ),
        (
            "Conclusions and Recommendations",
            "Strategic priorities for the upcoming product launch cycle"
        ),
    ]

    for i, (title_text, subtitle_text) in enumerate(slides_data):
        # Use Title Slide layout (index 0) for first slide, Title+Content (index 1) for rest
        layout_idx = 0 if i == 0 else 1
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Set title
        title_shape = slide.shapes.title
        title_shape.text = ""
        title_run = title_shape.text_frame.paragraphs[0].add_run()
        title_run.text = title_text
        title_run.font.bold = True
        title_run.font.size = Pt(32) if i == 0 else Pt(28)
        title_run.font.italic = False  # explicitly non-italic
        title_run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6D)

        # Set subtitle - placeholder index 1
        if 1 in [ph.placeholder_format.idx for ph in slide.placeholders]:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = ""
            sub_run = subtitle_shape.text_frame.paragraphs[0].add_run()
            sub_run.text = subtitle_text
            sub_run.font.size = Pt(18) if i == 0 else Pt(16)
            sub_run.font.italic = False  # explicitly non-italic
            sub_run.font.color.rgb = RGBColor(0x4A, 0x4A, 0x4A)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
