"""
Initial Setup: Create a 10-slide Project Plan presentation with an empty slide 8 titled 'Project Timeline'.
Task ID: impress_exec_074
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_074'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
    return slide


def add_title_only_slide(prs, title_text):
    """Add a slide with just a title and no body content."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title as a text box at top
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2B, 0x47, 0x9D)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(prs, "Project Plan: CloudSync Platform",
                    "Strategic Initiative 2025 | Prepared by Sarah Chen, VP Engineering")

    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "CloudSync Platform aims to unify our data synchronization layer across 14 regional data centers",
        "Expected 35% reduction in data latency and 22% cost savings on infrastructure",
        "Total project budget: $4.2M over 9 months with a team of 28 engineers",
        "Key stakeholders: Engineering, Product, DevOps, Security, and Customer Success teams",
        "Projected ROI of 280% within 18 months of launch",
    ])

    # Slide 3: Project Scope
    add_content_slide(prs, "Project Scope & Deliverables", [
        "Phase 1: Discovery and requirements gathering across 6 business units",
        "Phase 2: System architecture design with multi-region failover support",
        "Phase 3: Core platform development including API gateway and sync engine",
        "Phase 4: Comprehensive testing - unit, integration, load, and security testing",
        "Phase 5: Staged rollout starting with APAC region, then EMEA, then Americas",
    ])

    # Slide 4: Team Structure
    add_content_slide(prs, "Team Structure", [
        "Project Lead: Sarah Chen - Overall delivery and stakeholder management",
        "Tech Lead: Marcus Johnson - Architecture decisions and code quality",
        "Backend Team (8 engineers): Sync engine, data pipeline, API development",
        "Frontend Team (5 engineers): Dashboard, monitoring UI, admin console",
        "DevOps Team (4 engineers): CI/CD, infrastructure, deployment automation",
        "QA Team (6 engineers): Test automation, performance testing, security audits",
        "Product Manager: Elena Rodriguez - Feature prioritization and user research",
    ])

    # Slide 5: Budget Overview
    add_content_slide(prs, "Budget Allocation", [
        "Engineering salaries and contractors: $2,850,000 (67.8%)",
        "Cloud infrastructure and services: $620,000 (14.8%)",
        "Software licenses and tools: $340,000 (8.1%)",
        "Training and team development: $180,000 (4.3%)",
        "Contingency reserve: $210,000 (5.0%)",
        "Total approved budget: $4,200,000",
    ])

    # Slide 6: Risk Assessment
    add_content_slide(prs, "Risk Assessment", [
        "HIGH: Data migration complexity - Mitigation: staged migration with rollback",
        "HIGH: Third-party API dependency on Kafka 3.x - Mitigation: abstraction layer",
        "MEDIUM: Team scaling challenges - Mitigation: early hiring pipeline activation",
        "MEDIUM: Security compliance for PCI-DSS and SOC2 - Mitigation: dedicated security sprint",
        "LOW: Scope creep from stakeholder requests - Mitigation: change control board",
    ])

    # Slide 7: Key Milestones
    add_content_slide(prs, "Key Milestones", [
        "Jan 15: Project kickoff and team onboarding complete",
        "Feb 28: Discovery phase complete - requirements document signed off",
        "Apr 15: Architecture design approved by technical review board",
        "Jul 31: Core platform development complete - all APIs functional",
        "Aug 30: QA sign-off - all test suites passing with >95% coverage",
        "Sep 15: Production launch - APAC region live",
    ])

    # Slide 8: Project Timeline (EMPTY - task target)
    add_title_only_slide(prs, "Project Timeline")

    # Slide 9: Communication Plan
    add_content_slide(prs, "Communication Plan", [
        "Weekly standup: Monday 9:30 AM - All team leads (30 min)",
        "Bi-weekly sprint review: Friday 2:00 PM - Full team demo (60 min)",
        "Monthly steering committee: First Wednesday - Executive stakeholders (45 min)",
        "Quarterly business review: End of quarter - Board presentation (90 min)",
        "Ad-hoc escalation channel: #cloudsync-urgent Slack channel",
    ])

    # Slide 10: Next Steps
    add_content_slide(prs, "Immediate Next Steps", [
        "Finalize team allocation with department heads by Jan 10",
        "Set up development environments and CI/CD pipeline by Jan 12",
        "Schedule kickoff meetings with all 6 business units by Jan 14",
        "Complete initial technology stack evaluation by Jan 18",
        "Begin user research interviews with top 20 enterprise customers by Jan 20",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
