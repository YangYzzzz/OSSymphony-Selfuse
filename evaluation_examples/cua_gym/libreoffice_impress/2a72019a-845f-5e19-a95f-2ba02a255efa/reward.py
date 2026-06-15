"""
Reward Script: Gantt chart visualization on slide 8
Task ID: impress_exec_074
Domain: libreoffice_impress
Scoring:
  C1: Month labels (Jan-Sep) present on slide 8          — 0.20
  C2: Phase labels (5 phases) present on slide 8         — 0.15
  C3: 5 horizontal bar shapes exist on slide 8           — 0.25
  C4: All bar heights approximately 0.6 inches           — 0.15
  C5: Bar colors match phase specifications              — 0.25
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_074'

# Expected months across top
EXPECTED_MONTHS = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep']

# Expected phases and their approximate color families
# We use broad color family matching to be fair to reasonable variations
EXPECTED_PHASES = ['Discovery', 'Design', 'Development', 'Testing', 'Launch']

# Color family checks: (phase_name, color_family_description, checker_function)
def is_blue(rgb_str):
    """Blue family: high blue, low-mid red, low-mid green"""
    r, g, b = int(rgb_str[0:2], 16), int(rgb_str[2:4], 16), int(rgb_str[4:6], 16)
    return b > 150 and b > r and b > g

def is_green(rgb_str):
    """Green family: high green, lower red and blue"""
    r, g, b = int(rgb_str[0:2], 16), int(rgb_str[2:4], 16), int(rgb_str[4:6], 16)
    return g > 120 and g > r and g > b

def is_orange(rgb_str):
    """Orange family: high red, medium green, low blue"""
    r, g, b = int(rgb_str[0:2], 16), int(rgb_str[2:4], 16), int(rgb_str[4:6], 16)
    return r > 180 and g > 80 and g < 200 and b < 100

def is_purple(rgb_str):
    """Purple family: red and blue higher than green"""
    r, g, b = int(rgb_str[0:2], 16), int(rgb_str[2:4], 16), int(rgb_str[4:6], 16)
    return r > 80 and b > 80 and g < max(r, b) and (r > 100 or b > 100)

def is_red(rgb_str):
    """Red family: high red, low green and blue"""
    r, g, b = int(rgb_str[0:2], 16), int(rgb_str[2:4], 16), int(rgb_str[4:6], 16)
    return r > 150 and r > g and r > b and g < 130 and b < 130

PHASE_COLOR_CHECKS = {
    'Discovery': ('blue', is_blue),
    'Design': ('green', is_green),
    'Development': ('orange', is_orange),
    'Testing': ('purple', is_purple),
    'Launch': ('red', is_red),
}


def verify_task(file_path):
    """
    Verify Gantt chart creation on slide 8 with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 8 slides
    if len(prs.slides) < 8:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 8")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[7]  # slide 8, 0-indexed

    # Collect all text shapes and auto shapes on slide 8
    text_shapes = []
    auto_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                text_shapes.append((text, shape))
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            auto_shapes.append(shape)

    all_texts = [t for t, s in text_shapes]

    # -------------------------------------------------------
    # Component 1: Month labels (Jan-Sep) present (0.20 pts)
    # -------------------------------------------------------
    try:
        found_months = []
        for month in EXPECTED_MONTHS:
            # Check if month label exists as a standalone text or in any text shape
            if any(month.lower() == t.lower() for t, s in text_shapes):
                found_months.append(month)

        month_ratio = len(found_months) / len(EXPECTED_MONTHS)
        if month_ratio >= 0.8:  # at least 7 of 9 months
            pts = 0.20
            print(f"PASS: Component 1 — {len(found_months)}/{len(EXPECTED_MONTHS)} month labels found ({pts} pts)")
            total_score += pts
        elif month_ratio >= 0.5:
            pts = 0.10
            print(f"PARTIAL: Component 1 — {len(found_months)}/{len(EXPECTED_MONTHS)} month labels found ({pts} pts)")
            total_score += pts
        else:
            print(f"FAIL: Component 1 — Only {len(found_months)}/{len(EXPECTED_MONTHS)} month labels found: {found_months}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # -------------------------------------------------------
    # Component 2: Phase labels present (0.15 pts)
    # -------------------------------------------------------
    try:
        found_phases = []
        for phase in EXPECTED_PHASES:
            # Phase label can appear as standalone text box (left-side label)
            if any(phase.lower() == t.lower() for t, s in text_shapes):
                found_phases.append(phase)

        phase_ratio = len(found_phases) / len(EXPECTED_PHASES)
        if phase_ratio >= 0.8:  # at least 4 of 5
            pts = 0.15
            print(f"PASS: Component 2 — {len(found_phases)}/{len(EXPECTED_PHASES)} phase labels found ({pts} pts)")
            total_score += pts
        elif phase_ratio >= 0.5:
            pts = 0.075
            print(f"PARTIAL: Component 2 — {len(found_phases)}/{len(EXPECTED_PHASES)} phase labels found ({pts} pts)")
            total_score += pts
        else:
            print(f"FAIL: Component 2 — Only {len(found_phases)}/{len(EXPECTED_PHASES)} phase labels: {found_phases}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # -------------------------------------------------------
    # Component 3: 5 bar shapes (auto shapes) exist (0.25 pts)
    # -------------------------------------------------------
    try:
        # Bar shapes are auto shapes (rounded rectangles or rectangles) that are wider than tall
        bar_shapes = [s for s in auto_shapes if s.width > s.height]

        if len(bar_shapes) >= 5:
            pts = 0.25
            print(f"PASS: Component 3 — {len(bar_shapes)} horizontal bar shapes found ({pts} pts)")
            total_score += pts
        elif len(bar_shapes) >= 3:
            pts = 0.125
            print(f"PARTIAL: Component 3 — {len(bar_shapes)} horizontal bar shapes found ({pts} pts)")
            total_score += pts
        else:
            print(f"FAIL: Component 3 — Only {len(bar_shapes)} horizontal bar shapes found, need 5")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # -------------------------------------------------------
    # Component 4: Bar heights approximately 0.6 inches (0.15 pts)
    # -------------------------------------------------------
    try:
        bar_shapes = [s for s in auto_shapes if s.width > s.height]
        target_height = Inches(0.6)  # 548640 EMU
        tolerance = 0.15  # 15% tolerance

        correct_height_count = 0
        for bar in bar_shapes[:5]:  # check up to 5 bars
            actual_h = bar.height
            if abs(actual_h - target_height) / target_height <= tolerance:
                correct_height_count += 1

        if correct_height_count >= 4 and len(bar_shapes) >= 5:
            pts = 0.15
            print(f"PASS: Component 4 — {correct_height_count}/{min(len(bar_shapes),5)} bars have correct height ~0.6in ({pts} pts)")
            total_score += pts
        elif correct_height_count >= 2:
            pts = 0.075
            print(f"PARTIAL: Component 4 — {correct_height_count}/{min(len(bar_shapes),5)} bars have correct height ({pts} pts)")
            total_score += pts
        else:
            print(f"FAIL: Component 4 — Only {correct_height_count} bars have correct height (~0.6in)")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # -------------------------------------------------------
    # Component 5: Bar colors match phase specs (0.25 pts)
    # -------------------------------------------------------
    try:
        bar_shapes = [s for s in auto_shapes if s.width > s.height]

        # Build mapping: phase_name -> bar shape by matching text content
        phase_bars = {}
        for bar in bar_shapes:
            if bar.has_text_frame:
                bar_text = bar.text_frame.text.strip()
                for phase in EXPECTED_PHASES:
                    if phase.lower() == bar_text.lower():
                        phase_bars[phase] = bar
                        break

        correct_colors = 0
        checked = 0
        for phase in EXPECTED_PHASES:
            if phase not in phase_bars:
                # Try matching by vertical position order if text matching fails
                continue
            bar = phase_bars[phase]
            checked += 1
            try:
                fill = bar.fill
                if fill.type is not None:
                    rgb_str = str(fill.fore_color.rgb)
                    color_name, checker = PHASE_COLOR_CHECKS[phase]
                    if checker(rgb_str):
                        correct_colors += 1
                        print(f"  Color OK: {phase} = {rgb_str} (expected {color_name})")
                    else:
                        print(f"  Color MISMATCH: {phase} = {rgb_str} (expected {color_name})")
                else:
                    print(f"  No solid fill for {phase} bar")
            except Exception as e:
                print(f"  Color check error for {phase}: {e}")

        if correct_colors >= 4 and checked >= 4:
            pts = 0.25
            print(f"PASS: Component 5 — {correct_colors}/{checked} bar colors correct ({pts} pts)")
            total_score += pts
        elif correct_colors >= 2:
            pts = 0.125
            print(f"PARTIAL: Component 5 — {correct_colors}/{checked} bar colors correct ({pts} pts)")
            total_score += pts
        else:
            print(f"FAIL: Component 5 — Only {correct_colors}/{checked} bar colors correct")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
