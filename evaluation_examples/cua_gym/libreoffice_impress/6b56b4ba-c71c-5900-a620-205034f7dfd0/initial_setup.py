"""
Initial Setup: Create a 30-slide Lecture Notes presentation for print settings task.
Task ID: impress_el_014
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_el_014'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard slide dimensions (widescreen 10x7.5)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Lecture topics for 30 slides across a computer science course
    lecture_topics = [
        ("Introduction to Computer Science", "Course Overview and Objectives"),
        ("History of Computing", "From Babbage to Modern Processors"),
        ("Number Systems", "Binary, Octal, Hexadecimal Representations"),
        ("Boolean Algebra", "Logic Gates and Truth Tables"),
        ("Computer Architecture", "CPU, Memory, and I/O Systems"),
        ("Assembly Language", "Low-Level Programming Concepts"),
        ("Operating Systems", "Process Management and Scheduling"),
        ("Memory Management", "Virtual Memory and Paging"),
        ("File Systems", "Storage Organization and Access Methods"),
        ("Data Structures: Arrays", "Sequential Storage and Operations"),
        ("Data Structures: Linked Lists", "Dynamic Memory Allocation"),
        ("Data Structures: Stacks and Queues", "LIFO and FIFO Principles"),
        ("Data Structures: Trees", "Binary Trees and Traversals"),
        ("Data Structures: Graphs", "Representations and Applications"),
        ("Sorting Algorithms", "Bubble, Merge, and Quick Sort Analysis"),
        ("Searching Algorithms", "Linear and Binary Search Techniques"),
        ("Algorithm Complexity", "Big-O Notation and Performance Analysis"),
        ("Recursion", "Recursive Problem Solving Strategies"),
        ("Object-Oriented Programming", "Classes, Inheritance, Polymorphism"),
        ("Database Systems", "Relational Models and SQL Fundamentals"),
        ("Networking Fundamentals", "TCP/IP Protocol Stack and Layers"),
        ("Web Technologies", "HTTP, HTML, CSS, and JavaScript Basics"),
        ("Software Engineering", "SDLC Models and Agile Methodology"),
        ("Version Control", "Git Workflows and Collaboration"),
        ("Testing and Debugging", "Unit Tests, Integration Tests, CI/CD"),
        ("Cybersecurity Basics", "Encryption, Authentication, Threats"),
        ("Artificial Intelligence", "Machine Learning and Neural Networks"),
        ("Cloud Computing", "Virtualization, Containers, Microservices"),
        ("Ethics in Computing", "Privacy, Bias, and Responsible AI"),
        ("Course Summary", "Key Takeaways and Final Exam Preparation"),
    ]

    # Detailed bullet content for each slide
    slide_bullets = [
        ["Welcome to CS 101 - Fall 2025", "Instructor: Dr. Elena Rodriguez", "Office Hours: Mon/Wed 2-4 PM, Room 312"],
        ["Charles Babbage's Analytical Engine (1837)", "ENIAC - First general-purpose computer (1945)", "Moore's Law and exponential growth"],
        ["Decimal to binary conversion methods", "Hexadecimal notation in memory addresses", "Two's complement for signed integers"],
        ["AND, OR, NOT, XOR gate operations", "De Morgan's theorems and simplification", "Circuit design from Boolean expressions"],
        ["Von Neumann vs Harvard architecture", "Instruction fetch-decode-execute cycle", "Cache hierarchy: L1, L2, L3"],
        ["Registers, opcodes, and addressing modes", "MIPS instruction set examples", "Comparing RISC vs CISC architectures"],
        ["Process states: new, ready, running, waiting", "Round-robin and priority scheduling", "Context switching overhead analysis"],
        ["Physical vs logical address spaces", "Page tables and translation lookaside buffer", "Thrashing prevention strategies"],
        ["FAT32, NTFS, ext4 comparison", "Inodes and directory structures in Unix", "Journaling for crash recovery"],
        ["Static vs dynamic arrays", "Time complexity: O(1) access, O(n) insertion", "Multi-dimensional arrays and matrix operations"],
        ["Singly, doubly, and circular variants", "Pointer manipulation and memory leaks", "Comparison with array-based lists"],
        ["Stack applications: expression evaluation", "Queue applications: BFS, print spooling", "Priority queues and heap implementation"],
        ["BST properties and balanced trees (AVL, Red-Black)", "Tree traversal: inorder, preorder, postorder", "Applications: file systems, expression parsing"],
        ["Adjacency matrix vs adjacency list", "BFS and DFS traversal algorithms", "Shortest path: Dijkstra's algorithm"],
        ["Comparison-based sorting lower bound: O(n log n)", "Stability and in-place properties", "Counting sort and radix sort for integers"],
        ["Binary search requires sorted data", "Hash tables: O(1) average lookup", "Collision resolution: chaining vs open addressing"],
        ["Best, worst, and average case analysis", "Space complexity considerations", "Amortized analysis for dynamic arrays"],
        ["Base case and recursive case identification", "Call stack and stack overflow risks", "Dynamic programming as optimization of recursion"],
        ["Encapsulation, abstraction, inheritance", "Design patterns: Singleton, Observer, Factory", "SOLID principles for maintainable code"],
        ["Entity-relationship diagrams", "Normal forms: 1NF, 2NF, 3NF, BCNF", "JOIN operations and query optimization"],
        ["OSI model vs TCP/IP model", "IP addressing and subnetting", "DNS resolution and routing protocols"],
        ["Client-server architecture", "RESTful API design principles", "Frontend frameworks and SPA architecture"],
        ["Waterfall vs Agile vs DevOps", "Requirements gathering and specification", "Code reviews and pair programming"],
        ["Repository management and branching", "Merge conflicts and resolution strategies", "Continuous integration pipelines"],
        ["Test-driven development methodology", "Code coverage metrics and tools", "Debugging techniques and profiling"],
        ["Symmetric vs asymmetric encryption", "OAuth 2.0 and multi-factor authentication", "Common vulnerabilities: SQL injection, XSS"],
        ["Supervised vs unsupervised learning", "Convolutional and recurrent neural networks", "Natural language processing applications"],
        ["IaaS, PaaS, SaaS service models", "Docker containers and Kubernetes orchestration", "Serverless computing and event-driven architecture"],
        ["Data privacy regulations (GDPR, CCPA)", "Algorithmic bias and fairness", "Environmental impact of computing"],
        ["Review of core CS fundamentals", "Exam format: 40 MCQ + 3 coding problems", "Study resources and practice materials"],
    ]

    for i, (title, subtitle) in enumerate(lecture_topics):
        if i == 0:
            # Title slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title
            slide.placeholders[1].text = subtitle
        else:
            # Content slides with title + bullets
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Lecture {i}: {title}"

            # Add bullet content
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            bullets = slide_bullets[i]
            for j, bullet_text in enumerate(bullets):
                if j == 0:
                    tf.paragraphs[0].text = bullet_text
                    tf.paragraphs[0].font.size = Pt(18)
                else:
                    p = tf.add_paragraph()
                    p.text = bullet_text
                    p.font.size = Pt(18)
                    p.level = 0

        # Add slide number in notes for reference
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = f"Slide {i + 1} of 30 - CS 101 Fall 2025"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=3.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
