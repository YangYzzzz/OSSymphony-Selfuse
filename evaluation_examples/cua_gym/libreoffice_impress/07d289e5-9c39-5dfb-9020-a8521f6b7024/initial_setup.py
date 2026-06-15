"""
Initial Setup: Create Shape_Parade presentation with 4 shapes on slide 2 (no animations)
Task ID: impress_ma_075
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import math

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_075'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_freeform_triangle(slide, left, top, width, height):
    """Add a triangle as a freeform shape."""
    # Use an isosceles triangle auto shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        left, top, width, height
    )
    return shape


def add_freeform_star(slide, left, top, width, height):
    """Add a 5-pointed star auto shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.STAR_5_POINT,
        left, top, width, height
    )
    return shape


def create_initial():
    prs = Presentation()
    # Standard 16:9 slide size
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Shape Parade"
    slide1.placeholders[1].text = "Interactive Animation Showcase"

    # --- Slide 2: Four shapes in a row (NO animations) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide2.shapes.title.text = "Our Shape Collection"

    shape_size = Inches(1.8)
    y_pos = Inches(3.0)
    spacing = Inches(2.6)
    start_x = Inches(1.5)

    # Circle (oval)
    circle = slide2.shapes.add_shape(
        MSO_SHAPE.OVAL,
        start_x, y_pos, shape_size, shape_size
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x41, 0x72, 0xC4)  # Blue
    circle.name = "Circle"
    # Add label
    tf = circle.text_frame
    tf.paragraphs[0].text = "Circle"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for run in tf.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.size = Pt(14)
        run.font.bold = True

    # Square (rectangle)
    square = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        start_x + spacing, y_pos, shape_size, shape_size
    )
    square.fill.solid()
    square.fill.fore_color.rgb = RGBColor(0xED, 0x7D, 0x31)  # Orange
    square.name = "Square"
    tf = square.text_frame
    tf.paragraphs[0].text = "Square"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for run in tf.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.size = Pt(14)
        run.font.bold = True

    # Triangle
    triangle = add_freeform_triangle(
        slide2,
        start_x + spacing * 2, y_pos, shape_size, shape_size
    )
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = RGBColor(0x70, 0xAD, 0x47)  # Green
    triangle.name = "Triangle"
    tf = triangle.text_frame
    tf.paragraphs[0].text = "Triangle"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for run in tf.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.size = Pt(14)
        run.font.bold = True

    # Star
    star = add_freeform_star(
        slide2,
        start_x + spacing * 3, y_pos, shape_size, shape_size
    )
    star.fill.solid()
    star.fill.fore_color.rgb = RGBColor(0xFF, 0xC0, 0x00)  # Gold
    star.name = "Star"
    tf = star.text_frame
    tf.paragraphs[0].text = "Star"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for run in tf.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.size = Pt(14)
        run.font.bold = True

    # --- Slide 3: Project Timeline ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide3.shapes.title.text = "Project Timeline"
    body3 = slide3.placeholders[1]
    tf3 = body3.text_frame
    tf3.text = "Q1 2025: Initial concept development and prototyping"
    p2 = tf3.add_paragraph()
    p2.text = "Q2 2025: User testing and feedback collection"
    p3 = tf3.add_paragraph()
    p3.text = "Q3 2025: Final design refinements and launch preparation"
    p4 = tf3.add_paragraph()
    p4.text = "Q4 2025: Public release and marketing campaign"

    # --- Slide 4: Team Members ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Team Members"
    body4 = slide4.placeholders[1]
    tf4 = body4.text_frame
    tf4.text = "Elena Rodriguez - Creative Director"
    items = [
        "James Liu - Lead Developer",
        "Priya Sharma - UX Research Lead",
        "Marcus Webb - Marketing Strategist",
        "Aiko Tanaka - Quality Assurance Manager"
    ]
    for item in items:
        p = tf4.add_paragraph()
        p.text = item

    # --- Slide 5: Thank You ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[0])
    slide5.shapes.title.text = "Thank You!"
    slide5.placeholders[1].text = "Questions & Discussion"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
