"""
Reward Script: Underline all text in the content textbox on slide 2.
Task ID: osworld_impress_underline_darkred_table_001
Domain: libreoffice_impress
Scoring:
  - Component 1: Bullet point 1 in content textbox on slide 2 is underlined (0.33 pts)
  - Component 2: Bullet point 2 in content textbox on slide 2 is underlined (0.34 pts)
  - Component 3: Bullet point 3 in content textbox on slide 2 is underlined (0.33 pts)
Total: 1.0
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_underline_darkred_table_001'


def get_all_text_shapes(slide):
    """Recursively get all shapes with text frames, including inside groups."""
    def extract(shape):
        results = []
        if hasattr(shape, "text_frame"):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.extend(extract(sub))
        return results
    out = []
    for shape in slide.shapes:
        out.extend(extract(shape))
    return out


def check_run_underline(run):
    """Return True if the run has underline set to True."""
    return run.font.underline is True


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Task: Underline all text in the content textbox on slide 2.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: verify we have at least 2 slides
    if len(prs.slides) < 2:
        print(f"CRITICAL: Expected at least 2 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Get slide 2 (index 1)
    slide2 = prs.slides[1]

    # Find the content textbox on slide 2 (not the title)
    # The content textbox is named "Content Placeholder 2"
    content_shape = None
    for shape in get_all_text_shapes(slide2):
        # Skip title shapes; find the body/content placeholder
        if shape.has_text_frame and 'Content' in shape.name:
            content_shape = shape
            break

    if content_shape is None:
        # Fallback: find any non-title textbox on slide 2
        for shape in get_all_text_shapes(slide2):
            if shape.has_text_frame and shape.name != 'Title 1':
                content_shape = shape
                break

    if content_shape is None:
        print("CRITICAL: Cannot find content textbox on slide 2")
        print("REWARD: 0.0")
        return 0.0

    paragraphs = content_shape.text_frame.paragraphs
    # Filter to non-empty paragraphs
    nonempty_paras = [p for p in paragraphs if p.text.strip()]

    if len(nonempty_paras) == 0:
        print("CRITICAL: Content textbox on slide 2 has no non-empty paragraphs")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Found content textbox '{content_shape.name}' with {len(nonempty_paras)} non-empty paragraphs")

    # Component 1: Bullet point 1 in content textbox is underlined (0.33 pts)
    try:
        para = nonempty_paras[0]
        nonempty_runs = [r for r in para.runs if (r.text or "").strip()]
        if not nonempty_runs:
            print(f"FAIL: Component 1 — Paragraph 1 has no runs with text: '{para.text}'")
        else:
            all_underlined = all(check_run_underline(r) for r in nonempty_runs)
            if all_underlined:
                print(f"PASS: Component 1 — Paragraph 1 is underlined: '{para.text[:50]}' (0.33 pts)")
                total_score += 0.33
            else:
                underlined_statuses = [r.font.underline for r in nonempty_runs]
                print(f"FAIL: Component 1 — Paragraph 1 not fully underlined. underline statuses: {underlined_statuses}, text: '{para.text[:50]}'")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Bullet point 2 in content textbox is underlined (0.34 pts)
    try:
        if len(nonempty_paras) < 2:
            print(f"FAIL: Component 2 — Only {len(nonempty_paras)} non-empty paragraph(s) found, expected at least 2")
        else:
            para = nonempty_paras[1]
            nonempty_runs = [r for r in para.runs if (r.text or "").strip()]
            if not nonempty_runs:
                print(f"FAIL: Component 2 — Paragraph 2 has no runs with text: '{para.text}'")
            else:
                all_underlined = all(check_run_underline(r) for r in nonempty_runs)
                if all_underlined:
                    print(f"PASS: Component 2 — Paragraph 2 is underlined: '{para.text[:50]}' (0.34 pts)")
                    total_score += 0.34
                else:
                    underlined_statuses = [r.font.underline for r in nonempty_runs]
                    print(f"FAIL: Component 2 — Paragraph 2 not fully underlined. underline statuses: {underlined_statuses}, text: '{para.text[:50]}'")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Bullet point 3 in content textbox is underlined (0.33 pts)
    try:
        if len(nonempty_paras) < 3:
            print(f"FAIL: Component 3 — Only {len(nonempty_paras)} non-empty paragraph(s) found, expected at least 3")
        else:
            para = nonempty_paras[2]
            nonempty_runs = [r for r in para.runs if (r.text or "").strip()]
            if not nonempty_runs:
                print(f"FAIL: Component 3 — Paragraph 3 has no runs with text: '{para.text}'")
            else:
                all_underlined = all(check_run_underline(r) for r in nonempty_runs)
                if all_underlined:
                    print(f"PASS: Component 3 — Paragraph 3 is underlined: '{para.text[:50]}' (0.33 pts)")
                    total_score += 0.33
                else:
                    underlined_statuses = [r.font.underline for r in nonempty_runs]
                    print(f"FAIL: Component 3 — Paragraph 3 not fully underlined. underline statuses: {underlined_statuses}, text: '{para.text[:50]}'")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path on the VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
