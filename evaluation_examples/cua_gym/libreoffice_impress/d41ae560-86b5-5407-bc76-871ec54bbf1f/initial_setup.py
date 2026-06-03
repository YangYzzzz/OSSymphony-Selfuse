"""
Initial Setup: Underline all text in the content textbox on slide 2.
Task ID: osworld_impress_underline_darkred_table_001
Domain: libreoffice_impress

Creates a 5-slide business summary deck. Slide 2 has a body textbox with
3 bullet points in black text (no underline formatting applied yet).
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_underline_darkred_table_001'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # Use standard 16:9 dimensions
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q1 Business Summary"
    slide1.placeholders[1].text = "Strategic Overview · January–March 2025"

    # ---- Slide 2: Content slide with 3 bullet points (NO underline) ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Key Highlights"

    # Body text box (placeholder index 1 is the content area)
    body = slide2.placeholders[1]
    tf = body.text_frame
    tf.clear()

    # Bullet 1
    p1 = tf.paragraphs[0]
    p1.text = "Revenue increased by 18% compared to Q4 last year"
    r1 = p1.runs[0]
    r1.font.size = Pt(20)
    r1.font.bold = False
    r1.font.italic = False
    r1.font.underline = False
    r1.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # Bullet 2
    p2 = tf.add_paragraph()
    p2.text = "New client acquisitions reached a record 42 accounts this quarter"
    r2 = p2.runs[0]
    r2.font.size = Pt(20)
    r2.font.bold = False
    r2.font.italic = False
    r2.font.underline = False
    r2.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # Bullet 3
    p3 = tf.add_paragraph()
    p3.text = "Operating costs reduced by 9% following the restructuring initiative"
    r3 = p3.runs[0]
    r3.font.size = Pt(20)
    r3.font.bold = False
    r3.font.italic = False
    r3.font.underline = False
    r3.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # ---- Slide 3: Regional Performance ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Regional Performance"
    body3 = slide3.placeholders[1]
    tf3 = body3.text_frame
    tf3.clear()

    p3_1 = tf3.paragraphs[0]
    p3_1.text = "APAC region contributed 34% of total revenue — up from 28%"
    p3_1.runs[0].font.size = Pt(20)
    p3_1.runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    p3_2 = tf3.add_paragraph()
    p3_2.text = "EMEA maintained steady growth at 12% year-over-year"
    p3_2.runs[0].font.size = Pt(20)
    p3_2.runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    p3_3 = tf3.add_paragraph()
    p3_3.text = "Americas leads total volume with $3.2M in closed deals"
    p3_3.runs[0].font.size = Pt(20)
    p3_3.runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # ---- Slide 4: Product Updates ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Product Updates"
    body4 = slide4.placeholders[1]
    tf4 = body4.text_frame
    tf4.clear()

    p4_1 = tf4.paragraphs[0]
    p4_1.text = "Platform v4.2 launched with advanced analytics dashboard"
    p4_1.runs[0].font.size = Pt(20)
    p4_1.runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    p4_2 = tf4.add_paragraph()
    p4_2.text = "Mobile application downloads surpassed 500K milestone"
    p4_2.runs[0].font.size = Pt(20)
    p4_2.runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    p4_3 = tf4.add_paragraph()
    p4_3.text = "API integrations with 15 new enterprise partners completed"
    p4_3.runs[0].font.size = Pt(20)
    p4_3.runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # ---- Slide 5: Outlook ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Q2 Outlook & Priorities"
    body5 = slide5.placeholders[1]
    tf5 = body5.text_frame
    tf5.clear()

    p5_1 = tf5.paragraphs[0]
    p5_1.text = "Target 22% revenue growth with expanded sales team of 30"
    p5_1.runs[0].font.size = Pt(20)
    p5_1.runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    p5_2 = tf5.add_paragraph()
    p5_2.text = "Launch customer loyalty program in May across all regions"
    p5_2.runs[0].font.size = Pt(20)
    p5_2.runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    p5_3 = tf5.add_paragraph()
    p5_3.text = "Complete migration to cloud infrastructure by end of June"
    p5_3.runs[0].font.size = Pt(20)
    p5_3.runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open the initial artifact in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
