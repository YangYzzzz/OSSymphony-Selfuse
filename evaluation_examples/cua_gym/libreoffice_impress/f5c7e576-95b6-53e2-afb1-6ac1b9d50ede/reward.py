"""
Reward Script: Apply Calibri font and 18pt size to all text on slides 2, 3, and 4.
Task ID: osworld_impress_global_font_change_006
Domain: libreoffice_impress
Scoring:
  Component 1 (0.33): All non-empty runs on Slide 2 use Calibri font at 18pt
  Component 2 (0.33): All non-empty runs on Slide 3 use Calibri font at 18pt
  Component 3 (0.34): All non-empty runs on Slide 4 use Calibri font at 18pt
  Total: 1.0
"""

import os
from pptx import Presentation
from pptx.util import Pt

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_global_font_change_006'

# Target font and size
TARGET_FONT = 'Calibri'
TARGET_SIZE_PT = 18
TARGET_SIZE_EMU = int(TARGET_SIZE_PT * 12700)  # 228600 EMU

def get_nonempty_runs_on_slide(slide):
    """Return list of (shape_name, run) for all non-empty runs in non-placeholder text shapes."""
    results = []
    for shape in slide.shapes:
        if shape.has_text_frame and not shape.is_placeholder:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        results.append((shape.name, run))
    return results


def verify_slide_font(prs, slide_idx, slide_label, expected_font, expected_size_emu):
    """
    Verify all non-empty text runs in non-placeholder shapes on a slide
    have the expected font name and size.
    Returns True if all runs pass, False otherwise.
    """
    slide = prs.slides[slide_idx]
    runs = get_nonempty_runs_on_slide(slide)

    if not runs:
        print(f"WARN: {slide_label} — no non-empty runs found in non-placeholder shapes")
        # No runs to check — cannot award points if there's nothing to verify
        return False

    failure_count = 0
    for shape_name, run in runs:
        font_ok = (run.font.name == expected_font)
        size_ok = (run.font.size == expected_size_emu)

        if not font_ok:
            print(f"FAIL: {slide_label} shape '{shape_name}' run='{run.text[:30]}': "
                  f"expected font '{expected_font}', found '{run.font.name}'")
            failure_count += 1
        if not size_ok:
            actual_pt = run.font.size / 12700 if run.font.size else None
            print(f"FAIL: {slide_label} shape '{shape_name}' run='{run.text[:30]}': "
                  f"expected size {TARGET_SIZE_PT}pt ({expected_size_emu} EMU), "
                  f"found {actual_pt}pt ({run.font.size} EMU)")
            failure_count += 1

    if failure_count == 0:
        print(f"PASS: {slide_label} — all {len(runs)} run(s) use '{expected_font}' at {TARGET_SIZE_PT}pt")

    return failure_count == 0


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0

    Scoring:
      Component 1 (0.33 pts): Slide 2 — all textbox runs changed to Calibri at 18pt
      Component 2 (0.33 pts): Slide 3 — all textbox runs changed to Calibri at 18pt
      Component 3 (0.34 pts): Slide 4 — all textbox runs changed to Calibri at 18pt
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Sanity check: must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"CRITICAL: Presentation has {len(prs.slides)} slides, expected at least 4")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide 2 (index 1) — all textbox runs use Calibri at 18pt (0.33 pts)
    try:
        slide2_ok = verify_slide_font(prs, 1, "Slide 2", TARGET_FONT, TARGET_SIZE_EMU)
        if slide2_ok:
            total_score += 0.33
        else:
            print(f"FAIL: Component 1 — Slide 2 text not fully converted to {TARGET_FONT} at {TARGET_SIZE_PT}pt")
    except Exception as e:
        print(f"ERROR: Component 1 (Slide 2) — {e}")

    # Component 2: Slide 3 (index 2) — all textbox runs use Calibri at 18pt (0.33 pts)
    try:
        slide3_ok = verify_slide_font(prs, 2, "Slide 3", TARGET_FONT, TARGET_SIZE_EMU)
        if slide3_ok:
            total_score += 0.33
        else:
            print(f"FAIL: Component 2 — Slide 3 text not fully converted to {TARGET_FONT} at {TARGET_SIZE_PT}pt")
    except Exception as e:
        print(f"ERROR: Component 2 (Slide 3) — {e}")

    # Component 3: Slide 4 (index 3) — all textbox runs use Calibri at 18pt (0.34 pts)
    try:
        slide4_ok = verify_slide_font(prs, 3, "Slide 4", TARGET_FONT, TARGET_SIZE_EMU)
        if slide4_ok:
            total_score += 0.34
        else:
            print(f"FAIL: Component 3 — Slide 4 text not fully converted to {TARGET_FONT} at {TARGET_SIZE_PT}pt")
    except Exception as e:
        print(f"ERROR: Component 3 (Slide 4) — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point: test against canonical artifact path on VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
