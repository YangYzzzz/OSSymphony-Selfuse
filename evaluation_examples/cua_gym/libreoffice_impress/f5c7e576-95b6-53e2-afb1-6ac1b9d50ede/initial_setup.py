"""
Initial Setup: 7-slide marketing campaign pitch presentation
Task ID: osworld_impress_global_font_change_006
Domain: libreoffice_impress
NOTE: Slides 2-4 intentionally use various NON-Calibri fonts at NON-18pt sizes.
      The task is to change those to Calibri/18pt.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_global_font_change_006'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_name, font_size_pt,
                 bold=False, color_rgb=None, alignment=PP_ALIGN.LEFT):
    """Helper: add a text box with a single run to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    if color_rgb:
        run.font.color.rgb = color_rgb
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title slide ────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "NovaBrand 2025 Marketing Campaign"
    slide1.placeholders[1].text = "Driving Growth Through Innovation"
    for shape in slide1.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = "Arial"

    # ── Slide 2: Campaign Overview (uses Georgia / 24pt — NOT Calibri/18pt) ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_text_box(slide2,
                 Inches(0.5), Inches(0.3), Inches(12), Inches(1.2),
                 "Campaign Overview",
                 font_name="Georgia", font_size_pt=28, bold=True,
                 color_rgb=RGBColor(0x1F, 0x49, 0x7D),
                 alignment=PP_ALIGN.LEFT)
    add_text_box(slide2,
                 Inches(0.5), Inches(1.7), Inches(12), Inches(5.0),
                 (
                     "This year's campaign focuses on three pillars: brand awareness, "
                     "customer retention, and digital expansion. Our target demographic "
                     "spans millennials and Gen-Z consumers across North America and Europe. "
                     "Key channels include social media, influencer partnerships, and "
                     "programmatic advertising."
                 ),
                 font_name="Georgia", font_size_pt=22,
                 color_rgb=RGBColor(0x33, 0x33, 0x33),
                 alignment=PP_ALIGN.LEFT)

    # ── Slide 3: Target Audience (uses Times New Roman / 20pt) ───────────────
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide3,
                 Inches(0.5), Inches(0.3), Inches(12), Inches(1.2),
                 "Target Audience",
                 font_name="Times New Roman", font_size_pt=26, bold=True,
                 color_rgb=RGBColor(0x1F, 0x49, 0x7D),
                 alignment=PP_ALIGN.LEFT)
    add_text_box(slide3,
                 Inches(0.5), Inches(1.7), Inches(12), Inches(5.0),
                 (
                     "Primary segment: Adults aged 22–38 with household income above $60k. "
                     "Secondary segment: Small business owners in the tech and retail sectors. "
                     "Psychographics: Innovation-oriented, value-driven, digitally native. "
                     "Geographic focus: Tier-1 and Tier-2 cities in the US, UK, and Canada."
                 ),
                 font_name="Times New Roman", font_size_pt=20,
                 color_rgb=RGBColor(0x33, 0x33, 0x33),
                 alignment=PP_ALIGN.LEFT)

    # ── Slide 4: Budget Allocation (uses Verdana / 16pt) ────────────────────
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide4,
                 Inches(0.5), Inches(0.3), Inches(12), Inches(1.2),
                 "Budget Allocation",
                 font_name="Verdana", font_size_pt=24, bold=True,
                 color_rgb=RGBColor(0x1F, 0x49, 0x7D),
                 alignment=PP_ALIGN.LEFT)
    add_text_box(slide4,
                 Inches(0.5), Inches(1.7), Inches(12), Inches(5.0),
                 (
                     "Total marketing budget for FY2025: $4.2 million. "
                     "Digital advertising: 40% ($1.68M). "
                     "Content production and creative: 25% ($1.05M). "
                     "Events and trade shows: 15% ($630K). "
                     "Influencer and affiliate programs: 12% ($504K). "
                     "Market research and analytics: 8% ($336K)."
                 ),
                 font_name="Verdana", font_size_pt=16,
                 color_rgb=RGBColor(0x33, 0x33, 0x33),
                 alignment=PP_ALIGN.LEFT)

    # ── Slide 5: Key Performance Indicators ────────────────────────────────
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide5,
                 Inches(0.5), Inches(0.3), Inches(12), Inches(1.2),
                 "Key Performance Indicators",
                 font_name="Arial", font_size_pt=28, bold=True,
                 color_rgb=RGBColor(0x1F, 0x49, 0x7D),
                 alignment=PP_ALIGN.LEFT)
    add_text_box(slide5,
                 Inches(0.5), Inches(1.7), Inches(12), Inches(5.0),
                 (
                     "Brand awareness: +25% lift measured via consumer survey Q3 2025. "
                     "Website traffic: Target 5M unique visitors per month by Q4. "
                     "Conversion rate: Improve from 2.1% to 3.5% across digital properties. "
                     "Customer acquisition cost: Reduce by 18% vs. FY2024 baseline. "
                     "Net Promoter Score: Maintain above 62 throughout the year."
                 ),
                 font_name="Arial", font_size_pt=18,
                 color_rgb=RGBColor(0x33, 0x33, 0x33),
                 alignment=PP_ALIGN.LEFT)

    # ── Slide 6: Timeline & Milestones ─────────────────────────────────────
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide6,
                 Inches(0.5), Inches(0.3), Inches(12), Inches(1.2),
                 "Timeline & Milestones",
                 font_name="Arial", font_size_pt=28, bold=True,
                 color_rgb=RGBColor(0x1F, 0x49, 0x7D),
                 alignment=PP_ALIGN.LEFT)
    add_text_box(slide6,
                 Inches(0.5), Inches(1.7), Inches(12), Inches(5.0),
                 (
                     "Q1 2025: Brand refresh and campaign creative finalization. "
                     "Q2 2025: Soft launch — digital ads, influencer seeding, PR blitz. "
                     "Q3 2025: Full campaign rollout with major media buys. "
                     "Q4 2025: Retargeting push and holiday promotions. "
                     "December 2025: Full performance review and 2026 planning kickoff."
                 ),
                 font_name="Arial", font_size_pt=18,
                 color_rgb=RGBColor(0x33, 0x33, 0x33),
                 alignment=PP_ALIGN.LEFT)

    # ── Slide 7: Summary & Next Steps ──────────────────────────────────────
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide7,
                 Inches(0.5), Inches(0.3), Inches(12), Inches(1.2),
                 "Summary & Next Steps",
                 font_name="Arial", font_size_pt=28, bold=True,
                 color_rgb=RGBColor(0x1F, 0x49, 0x7D),
                 alignment=PP_ALIGN.LEFT)
    add_text_box(slide7,
                 Inches(0.5), Inches(1.7), Inches(12), Inches(5.0),
                 (
                     "The NovaBrand 2025 Campaign positions us for a record year. "
                     "Immediate actions: Finalize agency contracts by January 31. "
                     "Schedule kick-off workshop with all department heads for February 10. "
                     "Confirm media placements with Outbrain and Google by February 15. "
                     "Questions and open discussion welcome — let's align on priorities."
                 ),
                 font_name="Arial", font_size_pt=18,
                 color_rgb=RGBColor(0x33, 0x33, 0x33),
                 alignment=PP_ALIGN.LEFT)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
