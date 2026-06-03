"""
Initial Setup: Create CSAT_Pitch presentation with 8 slides. Slide 5 has title
'Customer Satisfaction' but empty content (no chart).
Task ID: impress_sales_075
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_075'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_body_slide(prs, layout_idx, title_text, body_lines):
    """Add a slide with title and bullet-point body content."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title_text
    # Find the body placeholder (index 1)
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def add_title_only_slide(prs, title_text):
    """Add a slide with title only (layout 5=Blank or 6=Title Only)."""
    # Use layout 5 (blank) to have a clean slide with just a title textbox
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "CSAT Pitch Deck"
    slide1.placeholders[1].text = "Quarterly Customer Satisfaction Review\nGreenfield Analytics — Q1 2025"

    # --- Slide 2: Revenue Overview ---
    add_title_body_slide(prs, 1, "Revenue Overview", [
        "Total Q1 Revenue: $3.42M (+12% YoY)",
        "Recurring Revenue: $2.18M (64% of total)",
        "New Business: $1.24M from 38 new accounts",
        "Average Deal Size: $32,600 (up from $28,100)",
        "Net Revenue Retention: 118%",
    ])

    # --- Slide 3: Sales Pipeline ---
    add_title_body_slide(prs, 1, "Sales Pipeline", [
        "Active Opportunities: 142 deals worth $8.7M",
        "Stage 3+ Pipeline: $4.2M (48% of total)",
        "Average Sales Cycle: 47 days (down from 52)",
        "Win Rate: 34% overall, 41% for enterprise segment",
        "Top Verticals: Healthcare (28%), FinTech (22%), SaaS (19%)",
    ])

    # --- Slide 4: Market Analysis ---
    add_title_body_slide(prs, 1, "Market Analysis", [
        "TAM grew to $14.8B in 2024 (Gartner estimate)",
        "Our market share: 2.3% (up from 1.8%)",
        "Key competitors: Zenith Corp (8.1%), Apex Solutions (5.4%)",
        "Emerging opportunity: Mid-market segment ($5M-50M revenue)",
        "Regulatory tailwinds: GDPR compliance driving demand",
    ])

    # --- Slide 5: Customer Satisfaction (EMPTY — no chart) ---
    slide5 = add_title_only_slide(prs, "Customer Satisfaction")
    # Deliberately left empty — the task is to add the donut chart here.

    # --- Slide 6: Product Performance ---
    add_title_body_slide(prs, 1, "Product Performance", [
        "Platform uptime: 99.97% (SLA target: 99.9%)",
        "Feature adoption rate: 73% within first 30 days",
        "Support tickets: 1,247 resolved (avg resolution: 4.2 hours)",
        "NPS Score: 62 (industry benchmark: 44)",
        "Mobile app downloads: 18,400 (+35% QoQ)",
    ])

    # --- Slide 7: Team Goals ---
    add_title_body_slide(prs, 1, "Team Goals — Q2 2025", [
        "Achieve $4.0M quarterly revenue (+17% growth target)",
        "Reduce churn to below 3.5% monthly",
        "Launch enterprise self-service portal by June 15",
        "Hire 8 additional CSMs for APAC expansion",
        "Improve CSAT score to 4.4/5.0 or higher",
    ])

    # --- Slide 8: Thank You ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[0])
    slide8.shapes.title.text = "Thank You"
    slide8.placeholders[1].text = "Questions & Discussion\nContact: insights@greenfieldanalytics.com"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
