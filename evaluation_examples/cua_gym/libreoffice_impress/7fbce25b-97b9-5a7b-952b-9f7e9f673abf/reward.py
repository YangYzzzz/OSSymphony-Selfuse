"""
Reward Script: Donut chart on slide 5 with customer satisfaction data
Task ID: impress_sales_075
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25) — Doughnut chart exists on slide 5
  Component 2 (0.20) — Chart title is 'Customer Satisfaction Score: 4.2/5.0'
  Component 3 (0.15) — Categories: Very Satisfied, Satisfied, Neutral, Unsatisfied
  Component 4 (0.20) — Data values: 45, 35, 15, 5
  Component 5 (0.20) — Point colors: #4CAF50, #8BC34A, #FFC107, #F44336
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_075'


def persist_app_state():
    """Save any unsaved LibreOffice edits before verification."""
    try:
        os.environ["DISPLAY"] = ":0"
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        import time
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"PRECONDITION FAIL: Expected at least 5 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide5 = prs.slides[4]  # 0-indexed

    # Find chart shape on slide 5
    chart_shape = None
    for shape in slide5.shapes:
        if shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Doughnut chart exists on slide 5 (0.25 points)
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            # DOUGHNUT chart type is -4120
            if chart.chart_type == -4120:
                print(f"PASS: Component 1 — Doughnut chart found on slide 5 (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 — Chart found but type is {chart.chart_type}, expected DOUGHNUT (-4120)")
        else:
            print(f"FAIL: Component 1 — No chart found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # If no chart, remaining components cannot pass
    if chart_shape is None:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    chart = chart_shape.chart

    # Component 2: Chart title is 'Customer Satisfaction Score: 4.2/5.0' (0.20 points)
    try:
        if chart.has_title:
            title_text = chart.chart_title.text_frame.text.strip()
            expected_title = "Customer Satisfaction Score: 4.2/5.0"
            if title_text == expected_title:
                print(f"PASS: Component 2 — Chart title matches exactly (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 2 — Title is '{title_text}', expected '{expected_title}'")
        else:
            print(f"FAIL: Component 2 — Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Categories match (0.15 points)
    try:
        expected_cats = ['Very Satisfied', 'Satisfied', 'Neutral', 'Unsatisfied']
        actual_cats = []
        for plot in chart.plots:
            actual_cats = list(plot.categories)
            break

        if actual_cats == expected_cats:
            print(f"PASS: Component 3 — Categories match exactly (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 3 — Categories are {actual_cats}, expected {expected_cats}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Data values match 45, 35, 15, 5 (0.20 points)
    try:
        expected_values = [45.0, 35.0, 15.0, 5.0]
        series = list(chart.series)
        if len(series) > 0:
            actual_values = list(series[0].values)
            if actual_values == expected_values:
                print(f"PASS: Component 4 — Data values match exactly (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 4 — Values are {actual_values}, expected {expected_values}")
        else:
            print(f"FAIL: Component 4 — No series found in chart")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Point colors match #4CAF50, #8BC34A, #FFC107, #F44336 (0.20 points)
    try:
        expected_colors = ['4CAF50', '8BC34A', 'FFC107', 'F44336']
        series = list(chart.series)
        if len(series) > 0:
            actual_colors = []
            color_match_count = 0
            for j in range(min(4, len(series[0].points))):
                point = series[0].points[j]
                try:
                    fill = point.format.fill
                    if fill.type is not None:
                        color_hex = str(fill.fore_color.rgb)
                        actual_colors.append(color_hex)
                        if color_hex.upper() == expected_colors[j].upper():
                            color_match_count += 1
                    else:
                        actual_colors.append('None')
                except Exception as ce:
                    actual_colors.append(f'error: {ce}')

            if color_match_count == 4:
                print(f"PASS: Component 5 — All 4 point colors match (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 5 — {color_match_count}/4 colors match. Actual: {actual_colors}, Expected: {expected_colors}")
        else:
            print(f"FAIL: Component 5 — No series found for color check")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
