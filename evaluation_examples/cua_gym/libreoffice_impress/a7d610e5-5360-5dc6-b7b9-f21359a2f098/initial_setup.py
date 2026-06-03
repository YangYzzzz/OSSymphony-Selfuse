"""
Initial Setup: Create a 15-slide lecture presentation with title but no navigation
Task ID: impress_teach_057
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_057'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=RGBColor(0x33, 0x33, 0x33), alignment=PP_ALIGN.LEFT):
    """Helper to add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    # Standard widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # =========== Slide 1: Title Slide (NO navigation) ===========
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Today's Lecture"
    slide1.placeholders[1].text = "Introduction to Data Structures and Algorithms"

    # =========== Slide 2: Agenda ===========
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "1. Arrays and Linked Lists"
    body2.add_paragraph().text = "2. Trees and Graph Traversal"
    body2.add_paragraph().text = "3. Sorting Algorithms"
    body2.add_paragraph().text = "4. Quiz Section"
    body2.add_paragraph().text = "5. Summary and Key Takeaways"

    # =========== Slides 3-5: Topic 1 - Arrays & Linked Lists ===========
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Topic 1: Arrays"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Arrays store elements in contiguous memory locations."
    body3.add_paragraph().text = "Access time: O(1) for index-based retrieval"
    body3.add_paragraph().text = "Insertion at end: O(1) amortized"
    body3.add_paragraph().text = "Insertion at arbitrary position: O(n)"
    body3.add_paragraph().text = "Memory overhead: minimal (just the elements)"

    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Topic 1: Linked Lists"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Singly linked list: each node points to next"
    body4.add_paragraph().text = "Doubly linked list: nodes point both forward and backward"
    body4.add_paragraph().text = "Insertion at head: O(1)"
    body4.add_paragraph().text = "Search: O(n) — must traverse from head"
    body4.add_paragraph().text = "Use cases: stacks, queues, undo history"

    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Topic 1: Comparison"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Arrays excel at random access and cache locality"
    body5.add_paragraph().text = "Linked lists excel at dynamic insertion/deletion"
    body5.add_paragraph().text = "Choose based on dominant operation pattern"
    body5.add_paragraph().text = "Hybrid: ArrayDeque combines benefits of both"

    # =========== Slides 6-8: Topic 2 - Trees & Graphs ===========
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Topic 2: Binary Trees"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Each node has at most two children"
    body6.add_paragraph().text = "Binary Search Tree (BST): left < parent < right"
    body6.add_paragraph().text = "Balanced BSTs: AVL trees, Red-Black trees"
    body6.add_paragraph().text = "Average search: O(log n), worst case: O(n)"

    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Topic 2: Graph Representations"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "Adjacency matrix: O(V²) space, O(1) edge lookup"
    body7.add_paragraph().text = "Adjacency list: O(V+E) space, efficient traversal"
    body7.add_paragraph().text = "Edge list: simple, good for Kruskal's algorithm"
    body7.add_paragraph().text = "Incidence matrix: used in network flow problems"

    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Topic 2: Graph Traversal"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "BFS: breadth-first, uses queue, finds shortest path"
    body8.add_paragraph().text = "DFS: depth-first, uses stack (or recursion)"
    body8.add_paragraph().text = "Dijkstra: weighted shortest path, O((V+E) log V)"
    body8.add_paragraph().text = "Topological sort: DAG ordering for dependencies"

    # =========== Slides 9-11: Topic 3 - Sorting ===========
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Topic 3: Comparison Sorts"
    body9 = slide9.placeholders[1].text_frame
    body9.text = "Bubble Sort: O(n²) — simple but slow"
    body9.add_paragraph().text = "Merge Sort: O(n log n) — stable, divide and conquer"
    body9.add_paragraph().text = "Quick Sort: O(n log n) average — in-place, fast in practice"
    body9.add_paragraph().text = "Heap Sort: O(n log n) — in-place, not stable"

    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "Topic 3: Non-Comparison Sorts"
    body10 = slide10.placeholders[1].text_frame
    body10.text = "Counting Sort: O(n+k) — integer keys in known range"
    body10.add_paragraph().text = "Radix Sort: O(d·(n+k)) — digit by digit"
    body10.add_paragraph().text = "Bucket Sort: O(n) average — uniformly distributed data"
    body10.add_paragraph().text = "Lower bound for comparison sorts: Ω(n log n)"

    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    slide11.shapes.title.text = "Topic 3: Choosing the Right Sort"
    body11 = slide11.placeholders[1].text_frame
    body11.text = "Small arrays (n < 50): Insertion sort is fastest"
    body11.add_paragraph().text = "General purpose: Timsort (Python/Java default)"
    body11.add_paragraph().text = "Memory constrained: Quick sort or Heap sort"
    body11.add_paragraph().text = "Stability needed: Merge sort or Timsort"

    # =========== Slide 12-13: Quiz ===========
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])
    slide12.shapes.title.text = "Quiz: Question 1"
    body12 = slide12.placeholders[1].text_frame
    body12.text = "What is the time complexity of searching in a balanced BST?"
    body12.add_paragraph().text = "A) O(1)"
    body12.add_paragraph().text = "B) O(log n)"
    body12.add_paragraph().text = "C) O(n)"
    body12.add_paragraph().text = "D) O(n log n)"

    slide13 = prs.slides.add_slide(prs.slide_layouts[1])
    slide13.shapes.title.text = "Quiz: Question 2"
    body13 = slide13.placeholders[1].text_frame
    body13.text = "Which sorting algorithm is NOT stable?"
    body13.add_paragraph().text = "A) Merge Sort"
    body13.add_paragraph().text = "B) Insertion Sort"
    body13.add_paragraph().text = "C) Heap Sort"
    body13.add_paragraph().text = "D) Bubble Sort"

    # =========== Slide 14: Summary ===========
    slide14 = prs.slides.add_slide(prs.slide_layouts[1])
    slide14.shapes.title.text = "Summary"
    body14 = slide14.placeholders[1].text_frame
    body14.text = "Arrays vs Linked Lists: choose by access pattern"
    body14.add_paragraph().text = "Trees enable efficient searching and hierarchical data"
    body14.add_paragraph().text = "Graphs model relationships; BFS/DFS are fundamental"
    body14.add_paragraph().text = "Sorting: know trade-offs between time, space, stability"
    body14.add_paragraph().text = "Practice implementing each data structure from scratch"

    # =========== Slide 15: References ===========
    slide15 = prs.slides.add_slide(prs.slide_layouts[1])
    slide15.shapes.title.text = "References & Further Reading"
    body15 = slide15.placeholders[1].text_frame
    body15.text = "Cormen et al. — Introduction to Algorithms (CLRS), 4th Ed."
    body15.add_paragraph().text = "Sedgewick & Wayne — Algorithms, 4th Edition"
    body15.add_paragraph().text = "Skiena — The Algorithm Design Manual, 3rd Edition"
    body15.add_paragraph().text = "LeetCode, HackerRank — Practice platforms"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
