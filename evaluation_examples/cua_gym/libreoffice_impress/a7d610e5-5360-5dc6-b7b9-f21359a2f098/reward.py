"""
Reward Script: Interactive lecture navigation system on slide 1
Task ID: impress_teach_057
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25) — Five rounded rectangle shapes on slide 1
  Component 2 (0.25) — All shapes have solid fill #1565C0
  Component 3 (0.25) — Correct labels with white 14pt bold text
  Component 4 (0.25) — Correct hyperlinks to specified slides
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_057'

# Expected button specs: (label, target_slide_filename)
EXPECTED_BUTTONS = {
    'Topic 1': 'slide3.xml',
    'Topic 2': 'slide6.xml',
    'Topic 3': 'slide9.xml',
    'Quiz': 'slide12.xml',
    'Summary': 'slide14.xml',
}


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.oxml.ns import qn
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 14 slides
    if len(prs.slides) < 14:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 14")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]

    # Find all auto shapes (rounded rectangles) on slide 1 that are NOT placeholders
    nav_shapes = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                from pptx.enum.shapes import MSO_SHAPE
                if shape.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
                    nav_shapes.append(shape)
            except Exception:
                # Fallback: check if name contains "Rounded Rectangle" or shape type ID 5
                nav_shapes.append(shape)

    # Component 1: Five rounded rectangle shapes exist on slide 1 (0.25 points)
    try:
        count = len(nav_shapes)
        if count == 5:
            print(f"PASS: Component 1 — Found exactly 5 rounded rectangle shapes on slide 1 (0.25 pts)")
            total_score += 0.25
        elif count >= 3:
            partial = 0.25 * (count / 5)
            print(f"PARTIAL: Component 1 — Found {count}/5 rounded rectangles ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — Found {count} rounded rectangles, expected 5")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    if len(nav_shapes) == 0:
        print("FAIL: No navigation shapes found. Cannot verify remaining components.")
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Build a mapping of shape text -> shape for found nav shapes
    shape_map = {}
    for shape in nav_shapes:
        text = ""
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text.strip()
        shape_map[text] = shape

    # Component 2: All shapes have solid fill #1565C0 (0.25 points)
    try:
        fill_pass = 0
        for label, shape in shape_map.items():
            try:
                fill = shape.fill
                if fill.type == 1:  # SOLID
                    rgb = str(fill.fore_color.rgb).upper()
                    if rgb == '1565C0':
                        fill_pass += 1
                    else:
                        print(f"  FAIL: Shape '{label}' fill color is {rgb}, expected 1565C0")
                else:
                    print(f"  FAIL: Shape '{label}' fill type is {fill.type}, expected SOLID (1)")
            except Exception as e:
                print(f"  ERROR: Shape '{label}' fill check: {e}")

        expected_count = min(len(shape_map), 5)
        if expected_count > 0 and fill_pass == expected_count:
            print(f"PASS: Component 2 — All {fill_pass} shapes have fill #1565C0 (0.25 pts)")
            total_score += 0.25
        elif fill_pass > 0:
            partial = 0.25 * (fill_pass / 5)
            print(f"PARTIAL: Component 2 — {fill_pass}/5 shapes have correct fill ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No shapes have correct fill #1565C0")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Correct labels with white 14pt bold text (0.25 points)
    try:
        text_pass = 0
        for expected_label in EXPECTED_BUTTONS.keys():
            if expected_label not in shape_map:
                print(f"  FAIL: No shape found with label '{expected_label}'")
                continue
            shape = shape_map[expected_label]
            # Check text formatting
            try:
                para = shape.text_frame.paragraphs[0]
                runs = [r for r in para.runs if (r.text or "").strip()]
                if not runs:
                    print(f"  FAIL: Shape '{expected_label}' has no text runs")
                    continue
                run = runs[0]
                # Check bold
                is_bold = run.font.bold is True
                # Check size (14pt = 177800 EMU)
                is_14pt = run.font.size is not None and abs(run.font.size - 177800) < 1000
                # Check white color
                is_white = False
                try:
                    if run.font.color.type is not None:
                        rgb = str(run.font.color.rgb).upper()
                        is_white = rgb == 'FFFFFF'
                except:
                    pass

                if is_bold and is_14pt and is_white:
                    text_pass += 1
                else:
                    reasons = []
                    if not is_bold:
                        reasons.append(f"bold={run.font.bold}")
                    if not is_14pt:
                        reasons.append(f"size={run.font.size}")
                    if not is_white:
                        try:
                            reasons.append(f"color={run.font.color.rgb}")
                        except:
                            reasons.append("color=None/theme")
                    print(f"  FAIL: Shape '{expected_label}' text formatting issues: {', '.join(reasons)}")
            except Exception as e:
                print(f"  ERROR: Shape '{expected_label}' text check: {e}")

        if text_pass == 5:
            print(f"PASS: Component 3 — All 5 labels have white 14pt bold text (0.25 pts)")
            total_score += 0.25
        elif text_pass > 0:
            partial = 0.25 * (text_pass / 5)
            print(f"PARTIAL: Component 3 — {text_pass}/5 labels have correct formatting ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No labels have correct text formatting")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Correct hyperlinks to specified slides (0.25 points)
    # Need to use XML to check hyperlink targets reliably
    try:
        link_pass = 0

        # Parse relationships from the pptx zip for slide 1
        rel_map = {}  # rId -> target
        with zipfile.ZipFile(file_path, 'r') as zf:
            try:
                with zf.open('ppt/slides/_rels/slide1.xml.rels') as f:
                    rel_tree = ET.parse(f)
                    rel_root = rel_tree.getroot()
                    ns_rel = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}
                    for rel in rel_root.findall('r:Relationship', ns_rel):
                        # Fallback: no namespace
                        rid = rel.get('Id')
                        target = rel.get('Target')
                        if rid and target:
                            rel_map[rid] = target
                    # Also try without namespace
                    if not rel_map:
                        for rel in rel_root.iter():
                            rid = rel.get('Id')
                            target = rel.get('Target')
                            if rid and target:
                                rel_map[rid] = target
            except KeyError:
                print("  ERROR: Could not find slide1 relationships file")

        # Now check each expected button's hyperlink
        for expected_label, expected_target in EXPECTED_BUTTONS.items():
            if expected_label not in shape_map:
                print(f"  FAIL: No shape with label '{expected_label}' for hyperlink check")
                continue
            shape = shape_map[expected_label]
            try:
                elem = shape._element
                # Find hlinkClick in cNvPr
                ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

                hlinkClick = None
                for cNvPr in elem.iter('{%s}cNvPr' % ns_p):
                    hlinkClick = cNvPr.find('{%s}hlinkClick' % ns_a)
                    if hlinkClick is not None:
                        break

                if hlinkClick is None:
                    print(f"  FAIL: Shape '{expected_label}' has no hyperlink")
                    continue

                rId = hlinkClick.get('{%s}id' % ns_r)
                actual_target = rel_map.get(rId, 'unknown')

                if actual_target == expected_target:
                    link_pass += 1
                else:
                    print(f"  FAIL: Shape '{expected_label}' links to '{actual_target}', expected '{expected_target}'")
            except Exception as e:
                print(f"  ERROR: Shape '{expected_label}' hyperlink check: {e}")

        if link_pass == 5:
            print(f"PASS: Component 4 — All 5 hyperlinks point to correct slides (0.25 pts)")
            total_score += 0.25
        elif link_pass > 0:
            partial = 0.25 * (link_pass / 5)
            print(f"PARTIAL: Component 4 — {link_pass}/5 hyperlinks correct ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — No hyperlinks point to correct slides")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
