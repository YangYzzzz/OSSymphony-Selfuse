"""
Reward Script: Create horizontal org chart on slide 4 with CEO and VP boxes
Task ID: impress_ndo_062
Domain: libreoffice_impress
Scoring:
  Component 1: Four rounded rectangles with correct text labels (0.30)
  Component 2: Correct fill colors - CEO #2C3E50, VPs #3498DB (0.20)
  Component 3: All text white 14pt (0.20)
  Component 4: Three elbow-style connectors present (0.20)
  Component 5: Horizontal layout - CEO left, VPs right (0.10)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_062'


def persist_app_state(domain):
    """Save any unsaved LibreOffice state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            import time
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)

    # Collect auto shapes (rounded rectangles) on slide 4
    auto_shapes = []
    for shape in slide.shapes:
        if shape.shape_type == 1:  # AUTO_SHAPE
            auto_shapes.append(shape)

    # Collect connectors on slide 4
    connectors = []
    for shape in slide.shapes:
        if shape.shape_type == 9:  # LINE/CONNECTOR
            connectors.append(shape)

    # =========================================================
    # Component 1: Four rounded rectangles with correct text (0.30 points)
    # =========================================================
    try:
        expected_labels = {'CEO', 'VP Sales', 'VP Engineering', 'VP Marketing'}
        found_labels = set()
        for shape in auto_shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text in expected_labels:
                    found_labels.add(text)

        if found_labels == expected_labels:
            print(f"PASS: Component 1 - All 4 required labels found: {found_labels} (0.30 pts)")
            total_score += 0.30
        else:
            missing = expected_labels - found_labels
            print(f"FAIL: Component 1 - Missing labels: {missing}. Found: {found_labels}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # =========================================================
    # Component 2: Correct fill colors (0.20 points)
    # CEO box #2C3E50, VP boxes #3498DB
    # =========================================================
    try:
        ceo_color_ok = False
        vp_color_count = 0
        vp_expected = 3

        for shape in auto_shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            try:
                fill = shape.fill
                if fill.type is not None and fill.type == 1:  # SOLID
                    rgb_val = str(fill.fore_color.rgb).upper()
                else:
                    rgb_val = None
            except Exception:
                rgb_val = None

            if text == 'CEO':
                if rgb_val and rgb_val == '2C3E50':
                    ceo_color_ok = True
                    print(f"  CEO fill: {rgb_val} == 2C3E50 OK")
                else:
                    print(f"  CEO fill: {rgb_val} != 2C3E50")
            elif text in {'VP Sales', 'VP Engineering', 'VP Marketing'}:
                if rgb_val and rgb_val == '3498DB':
                    vp_color_count += 1
                    print(f"  {text} fill: {rgb_val} == 3498DB OK")
                else:
                    print(f"  {text} fill: {rgb_val} != 3498DB")

        if ceo_color_ok and vp_color_count == vp_expected:
            print(f"PASS: Component 2 - All fill colors correct (0.20 pts)")
            total_score += 0.20
        elif ceo_color_ok or vp_color_count > 0:
            partial = 0.0
            if ceo_color_ok:
                partial += 0.05
            partial += 0.05 * vp_color_count
            print(f"PARTIAL: Component 2 - CEO={ceo_color_ok}, VP colors={vp_color_count}/{vp_expected} ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 - No correct fill colors found")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # =========================================================
    # Component 3: All text is white (FFFFFF) and 14pt (0.20 points)
    # =========================================================
    try:
        text_checks_total = 0
        text_checks_pass = 0

        for shape in auto_shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if text not in {'CEO', 'VP Sales', 'VP Engineering', 'VP Marketing'}:
                continue

            text_checks_total += 1
            color_ok = False
            size_ok = False

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    # Check color
                    try:
                        if run.font.color.type is not None:
                            rgb_str = str(run.font.color.rgb).upper()
                            if rgb_str == 'FFFFFF':
                                color_ok = True
                    except Exception:
                        pass

                    # Check size (14pt = 177800 EMU)
                    if run.font.size is not None:
                        size_pt = run.font.size / 12700
                        if abs(size_pt - 14.0) < 0.5:
                            size_ok = True

            if color_ok and size_ok:
                text_checks_pass += 1
                print(f"  {text}: white 14pt OK")
            else:
                print(f"  {text}: color_ok={color_ok}, size_ok={size_ok}")

        if text_checks_total > 0 and text_checks_pass == text_checks_total:
            print(f"PASS: Component 3 - All {text_checks_pass} shapes have white 14pt text (0.20 pts)")
            total_score += 0.20
        elif text_checks_pass > 0:
            partial = 0.20 * (text_checks_pass / max(text_checks_total, 1))
            print(f"PARTIAL: Component 3 - {text_checks_pass}/{text_checks_total} shapes correct ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 - No shapes with correct white 14pt text")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # =========================================================
    # Component 4: Three elbow-style connectors (0.20 points)
    # =========================================================
    try:
        elbow_connector_count = 0

        # Check via XML for bentConnector preset geometry
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        for shape in connectors:
            elem = shape._element
            geom_list = elem.findall(f'.//{{{ns_a}}}prstGeom')
            for geom in geom_list:
                prst = geom.get('prst', '')
                if 'bent' in prst.lower() or 'elbow' in prst.lower():
                    elbow_connector_count += 1
                    break

        if elbow_connector_count >= 3:
            print(f"PASS: Component 4 - Found {elbow_connector_count} elbow connectors (0.20 pts)")
            total_score += 0.20
        elif elbow_connector_count > 0:
            partial = 0.20 * (elbow_connector_count / 3)
            print(f"PARTIAL: Component 4 - Found {elbow_connector_count}/3 elbow connectors ({partial:.2f} pts)")
            total_score += round(partial, 2)
        else:
            # Check if there are any connectors at all (even non-elbow)
            if len(connectors) >= 3:
                print(f"PARTIAL: Component 4 - Found {len(connectors)} connectors but not elbow style (0.07 pts)")
                total_score += 0.07
            else:
                print(f"FAIL: Component 4 - Found {len(connectors)} connectors, need 3 elbow-style")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # =========================================================
    # Component 5: Horizontal layout - CEO on left, VPs on right (0.10 points)
    # =========================================================
    try:
        ceo_shape = None
        vp_shapes = []

        for shape in auto_shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if text == 'CEO':
                ceo_shape = shape
            elif text in {'VP Sales', 'VP Engineering', 'VP Marketing'}:
                vp_shapes.append(shape)

        if ceo_shape is not None and len(vp_shapes) >= 3:
            ceo_center_x = ceo_shape.left + ceo_shape.width // 2
            vp_centers_x = [(s.left + s.width // 2) for s in vp_shapes]

            # CEO should be to the left of all VP boxes
            all_right = all(vp_cx > ceo_center_x for vp_cx in vp_centers_x)

            if all_right:
                print(f"PASS: Component 5 - CEO (x_center={ceo_center_x}) is left of all VPs (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 5 - CEO center_x={ceo_center_x}, VP centers={vp_centers_x}")
        else:
            print(f"FAIL: Component 5 - Cannot verify layout (CEO found={ceo_shape is not None}, VP count={len(vp_shapes)})")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
