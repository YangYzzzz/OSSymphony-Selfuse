"""
Initial Setup: Create OrgChart presentation with empty slide 4 for org chart task
Task ID: impress_ndo_062
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_062'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "OrgChart"
    slide1.placeholders[1].text = "Company Organizational Structure\nFiscal Year 2025-2026"

    # --- Slide 2: Company Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Company Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Founded in 2018, Meridian Technologies has grown to over 450 employees across three continents."
    p2 = body2.add_paragraph()
    p2.text = "Headquarters: San Francisco, CA"
    p2.level = 1
    p3 = body2.add_paragraph()
    p3.text = "Revenue: $128M (2024)"
    p3.level = 1
    p4 = body2.add_paragraph()
    p4.text = "Key Markets: Enterprise SaaS, Cloud Infrastructure, Developer Tools"
    p4.level = 1

    # --- Slide 3: Strategic Priorities ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Strategic Priorities Q2 2025"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Expand into APAC markets with localized product offerings"
    for item in [
        "Launch next-generation analytics platform by June 2025",
        "Achieve SOC 2 Type II certification for all cloud services",
        "Grow engineering team by 35% through targeted recruitment",
        "Reduce customer churn rate from 8.2% to below 5%",
    ]:
        pp = body3.add_paragraph()
        pp.text = item
        pp.level = 0

    # --- Slide 4: Leadership Team (empty below title) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title text box manually on blank layout
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Leadership Team"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # Slide 4 is intentionally empty below the title - the agent must add the org chart

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
