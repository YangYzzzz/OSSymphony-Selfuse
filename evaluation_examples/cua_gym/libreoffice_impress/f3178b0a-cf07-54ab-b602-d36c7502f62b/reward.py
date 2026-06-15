"""
FINAL REWARD SCRIPT - SUCCESS
Task: I’m wrapping up a LibreOffice Impress file that currently ends at slide 19. How can I quickly clone slides 18 and 19 and drop those copies right after the existing last slide so they become the new slides 20 and 21?
Generated: 2025-09-10 12:58:59
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation

"""
Reward Script for Task:
"I’m wrapping up a LibreOffice Impress file that currently ends at slide 19. How can I quickly clone slides 18 and 19 and drop those copies right after the existing last slide so they become the new slides 20 and 21?"

Verification Logic:
1. Load the provided presentation file safely (no points for simply loading).
2. Requirement 1 (0.4 pts) – Slide count must now be 21 (original 19 + 2 cloned slides).
3. Requirement 2 (0.3 pts) – Slide 20 (index 19) must be an exact duplicate of original slide 18 (index 17) based on textual content.
4. Requirement 3 (0.3 pts) – Slide 21 (index 20) must be an exact duplicate of original slide 19 (index 18) based on textual content.

Progressive scoring ensures partial credit when only some conditions are met.
The script prints detailed diagnostics and finally prints "REWARD: X.X".
"""

# -------- Helper Functions --------

def _extract_slide_texts(slide):
    """Return list of cleaned text strings from all shapes in the slide."""
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            raw = shape.text
            if raw:
                cleaned = " ".join(raw.split())  # normalize whitespace
                if cleaned:
                    texts.append(cleaned)
    return texts

def _slides_equal(slide_a, slide_b):
    """Check if two slides are textually identical (order-sensitive)."""
    return _extract_slide_texts(slide_a) == _extract_slide_texts(slide_b)

# -------- Main Verification --------

def verify_clone_slides_task(file_path: str) -> float:
    """Verify that slides 18 & 19 were cloned to become new slides 20 & 21."""

    total_score = 0.0  # progressive score
    MAX_SCORE = 1.0

    # 1. Load presentation (no points for loading itself)
    try:
        prs = Presentation(file_path)
    except Exception as exc:
        print(f"✗ Failed to load presentation: {exc}")
        print("REWARD: 0.0")
        return 0.0

    slide_count = len(prs.slides)
    expected_count = 21  # 19 original + 2 clones
    print(f"Slide count found: {slide_count}, expected: {expected_count}")

    # 2. Verify slide count (0.4 pts)
    if slide_count == expected_count:
        total_score += 0.4
        print("✓ Slide count requirement satisfied (0.4)")
    else:
        print("✗ Slide count incorrect (0 pts)")

    # 3. Verify cloned slide 18 -> 20 (0.3 pts)
    # 4. Verify cloned slide 19 -> 21 (0.3 pts)
    if slide_count >= expected_count:
        original_slide_18 = prs.slides[17]  # index 17
        original_slide_19 = prs.slides[18]  # index 18
        new_slide_20      = prs.slides[19]  # index 19
        new_slide_21      = prs.slides[20]  # index 20

        if _slides_equal(original_slide_18, new_slide_20):
            total_score += 0.3
            print("✓ Slide 20 is an exact duplicate of Slide 18 (0.3)")
        else:
            print("✗ Slide 20 is NOT a duplicate of Slide 18 (0 pts)")

        if _slides_equal(original_slide_19, new_slide_21):
            total_score += 0.3
            print("✓ Slide 21 is an exact duplicate of Slide 19 (0.3)")
        else:
            print("✗ Slide 21 is NOT a duplicate of Slide 19 (0 pts)")
    else:
        print("✗ Not enough slides to perform duplication checks (0 pts)")

    # Cap score at 1.0
    final_score = min(total_score, MAX_SCORE)
    print(f"Total score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score

# ---------------- Run Verification ----------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/im_wrapping_up_a_libreoffice_impress_file_that_currently_ends_at_slide_19_how_can_i_quickly_clone_sl_golden.pptx"
    verify_clone_slides_task(FILE_PATH)
