"""
Reward Script: Apply global font change to entire presentation
Task ID: osworld_impress_global_font_change_011
Domain: libreoffice_impress
Scoring:
  Component 1: All text runs use Garamond font (0.4 pts)
  Component 2: All text runs are 16pt (203200 EMU) (0.3 pts)
  Component 3: All text runs are colored #003366 (0.3 pts)
  Total: 1.0
"""

import os
from pptx import Presentation
from pptx.util import Pt

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_global_font_change_011'

TARGET_FONT = 'Garamond'
TARGET_SIZE_EMU = 203200   # 16pt * 12700 EMU/pt
TARGET_COLOR = '003366'    # dark blue #003366


def persist_app_state():
    """Send Ctrl+S to save any unsaved GUI edits in LibreOffice Impress."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def get_all_runs(prs):
    """
    Iterate over every slide, every shape (including group sub-shapes),
    every paragraph, and every non-empty run, yielding (slide_idx, run) tuples.
    """
    def extract_shapes(shape):
        shapes = []
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                shapes.extend(extract_shapes(sub))
        else:
            shapes.append(shape)
        return shapes

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            for real_shape in extract_shapes(shape):
                if hasattr(real_shape, 'has_text_frame') and real_shape.has_text_frame:
                    for para in real_shape.text_frame.paragraphs:
                        for run in para.runs:
                            if (run.text or '').strip():
                                yield slide_idx, run


def verify_task(file_path):
    """
    Verify task completion: all text runs across all 9 slides must use
    Garamond font, 16pt size, and dark blue color #003366.
    Returns a float between 0.0 and 1.0.
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Collect all non-empty runs once
    all_runs = list(get_all_runs(prs))

    if not all_runs:
        print("CRITICAL: No text runs found in presentation.")
        print("REWARD: 0.0")
        return 0.0

    total_runs = len(all_runs)
    print(f"INFO: Found {total_runs} non-empty text runs across {len(prs.slides)} slides")

    # ------------------------------------------------------------------
    # Component 1: Font name is Garamond for ALL runs (0.4 points)
    # ------------------------------------------------------------------
    font_pass = 0
    font_fail_examples = []
    try:
        for slide_idx, run in all_runs:
            actual_font = run.font.name
            if actual_font == TARGET_FONT:
                font_pass += 1
            else:
                if len(font_fail_examples) < 3:
                    font_fail_examples.append(
                        f"slide {slide_idx + 1}: text={run.text[:20]!r}, font={actual_font!r}"
                    )

        if font_pass == total_runs:
            print(f"PASS: Component 1 — All {total_runs} runs use Garamond font (0.4 pts)")
            total_score += 0.4
        else:
            pct = font_pass / total_runs * 100
            print(f"FAIL: Component 1 — {font_pass}/{total_runs} runs ({pct:.0f}%) use Garamond.")
            if font_fail_examples:
                print(f"  First failures: {'; '.join(font_fail_examples)}")
    except Exception as e:
        print(f"ERROR: Component 1 (font name) — {e}")

    # ------------------------------------------------------------------
    # Component 2: Font size is 16pt (203200 EMU) for ALL runs (0.3 points)
    # ------------------------------------------------------------------
    size_pass = 0
    size_fail_examples = []
    try:
        for slide_idx, run in all_runs:
            actual_size = run.font.size
            if actual_size == TARGET_SIZE_EMU:
                size_pass += 1
            else:
                if len(size_fail_examples) < 3:
                    size_fail_examples.append(
                        f"slide {slide_idx + 1}: text={run.text[:20]!r}, size={actual_size}"
                    )

        if size_pass == total_runs:
            print(f"PASS: Component 2 — All {total_runs} runs are 16pt (203200 EMU) (0.3 pts)")
            total_score += 0.3
        else:
            pct = size_pass / total_runs * 100
            print(f"FAIL: Component 2 — {size_pass}/{total_runs} runs ({pct:.0f}%) are 16pt.")
            if size_fail_examples:
                print(f"  First failures: {'; '.join(size_fail_examples)}")
    except Exception as e:
        print(f"ERROR: Component 2 (font size) — {e}")

    # ------------------------------------------------------------------
    # Component 3: Font color is #003366 for ALL runs (0.3 points)
    # ------------------------------------------------------------------
    color_pass = 0
    color_fail_examples = []
    try:
        for slide_idx, run in all_runs:
            try:
                if run.font.color.type is not None:
                    actual_color = str(run.font.color.rgb).upper()
                    expected_color = TARGET_COLOR.upper()
                    if actual_color == expected_color:
                        color_pass += 1
                    else:
                        if len(color_fail_examples) < 3:
                            color_fail_examples.append(
                                f"slide {slide_idx + 1}: text={run.text[:20]!r}, color={actual_color}"
                            )
                else:
                    # No explicit color set — does not match #003366
                    if len(color_fail_examples) < 3:
                        color_fail_examples.append(
                            f"slide {slide_idx + 1}: text={run.text[:20]!r}, color=None (not set)"
                        )
            except Exception as inner_e:
                if len(color_fail_examples) < 3:
                    color_fail_examples.append(
                        f"slide {slide_idx + 1}: text={run.text[:20]!r}, color_err={inner_e}"
                    )

        if color_pass == total_runs:
            print(f"PASS: Component 3 — All {total_runs} runs have color #003366 (0.3 pts)")
            total_score += 0.3
        else:
            pct = color_pass / total_runs * 100
            print(f"FAIL: Component 3 — {color_pass}/{total_runs} runs ({pct:.0f}%) have color #003366.")
            if color_fail_examples:
                print(f"  First failures: {'; '.join(color_fail_examples)}")
    except Exception as e:
        print(f"ERROR: Component 3 (font color) — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path on the VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
