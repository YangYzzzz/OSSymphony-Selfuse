"""
Initial Setup: Apply a global font change to the entire presentation
Task ID: osworld_impress_global_font_change_011
Domain: libreoffice_impress

Creates a 9-slide annual report deck with mixed fonts and colors.
Each slide has at least a title and body text in various configurations.
DOES NOT contain Garamond font, 16pt size, or #003366 color (those are the task targets).
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_global_font_change_011'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_name, font_size_pt, bold=False, color_rgb=None, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    if color_rgb is not None:
        run.font.color.rgb = color_rgb
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Acme Corporation"
    slide1.shapes.title.text_frame.paragraphs[0].runs[0].font.name = "Arial"
    slide1.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(40)
    slide1.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x8B, 0x00, 0x00)

    subtitle_ph = slide1.placeholders[1]
    subtitle_ph.text = "Annual Report 2024 — A Year of Growth and Transformation"
    subtitle_ph.text_frame.paragraphs[0].runs[0].font.name = "Calibri"
    subtitle_ph.text_frame.paragraphs[0].runs[0].font.size = Pt(22)
    subtitle_ph.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ---- Slide 2: Executive Summary ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    slide2.shapes.title.text_frame.paragraphs[0].runs[0].font.name = "Times New Roman"
    slide2.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(36)
    slide2.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x47, 0xAB)

    body_ph = slide2.placeholders[1]
    body_ph.text = (
        "Total revenue reached $2.4 billion, a 12% increase year-over-year.\n"
        "Operating profit margin improved to 18.3% from 15.1% in 2023.\n"
        "Expanded into 7 new international markets across Southeast Asia.\n"
        "Launched 4 flagship product lines with cumulative sales of 830,000 units.\n"
        "Headcount grew from 4,200 to 5,100 employees globally."
    )
    for para in body_ph.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    # ---- Slide 3: Financial Highlights ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Financial Highlights"
    slide3.shapes.title.text_frame.paragraphs[0].runs[0].font.name = "Georgia"
    slide3.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(34)
    slide3.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1B, 0x5E, 0x20)

    body_ph3 = slide3.placeholders[1]
    body_ph3.text = (
        "Revenue: $2.4B  (+12% YoY)\n"
        "Gross Profit: $1.08B  (+14% YoY)\n"
        "Net Income: $312M  (+9% YoY)\n"
        "EBITDA: $480M  (+11% YoY)\n"
        "EPS: $4.28 (diluted)"
    )
    for para in body_ph3.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Verdana"
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x80)

    # Add a text box with additional detail
    add_text_box(
        slide3,
        Inches(8), Inches(3.5), Inches(4.5), Inches(2.5),
        "Q4 2024 highlighted record cash reserves of $620M,\nenabling accelerated share buybacks and dividends.",
        "Palatino Linotype", 14, bold=False,
        color_rgb=RGBColor(0x55, 0x55, 0x55),
        alignment=PP_ALIGN.LEFT
    )

    # ---- Slide 4: Regional Performance ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Regional Performance"
    slide4.shapes.title.text_frame.paragraphs[0].runs[0].font.name = "Franklin Gothic Medium"
    slide4.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    slide4.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x7B, 0x1F, 0xA2)

    body_ph4 = slide4.placeholders[1]
    body_ph4.text = (
        "North America: $960M  (40% of revenue)\n"
        "Europe & Middle East: $600M  (25%)\n"
        "Asia-Pacific: $480M  (20%)\n"
        "Latin America: $240M  (10%)\n"
        "Rest of World: $120M  (5%)"
    )
    for para in body_ph4.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Trebuchet MS"
            run.font.size = Pt(19)
            run.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    # ---- Slide 5: Product Portfolio ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Product Portfolio"
    slide5.shapes.title.text_frame.paragraphs[0].runs[0].font.name = "Cambria"
    slide5.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(35)
    slide5.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xBF, 0x36, 0x0C)

    body_ph5 = slide5.placeholders[1]
    body_ph5.text = (
        "ProSuite X1 — Enterprise software platform: 210,000 licenses sold\n"
        "CloudSync Pro — Data synchronization tool: 320,000 subscriptions\n"
        "DataGuard 360 — Cybersecurity suite: 180,000 deployments\n"
        "Insight Analytics — BI and reporting: 120,000 active users\n"
        "MobileFirst SDK — Developer toolkit: 1.2M downloads"
    )
    for para in body_ph5.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Gill Sans MT"
            run.font.size = Pt(17)
            run.font.color.rgb = RGBColor(0x01, 0x57, 0x9B)

    # ---- Slide 6: Human Capital ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Human Capital & Culture"
    slide6.shapes.title.text_frame.paragraphs[0].runs[0].font.name = "Garamond"
    # NOTE: Title uses Garamond intentionally on slide 6 — body deliberately uses a different font
    slide6.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(33)
    slide6.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xE6, 0x51, 0x00)

    body_ph6 = slide6.placeholders[1]
    body_ph6.text = (
        "5,100 employees across 28 countries\n"
        "Employee satisfaction score: 87% (up from 81%)\n"
        "Training hours per employee: 42 (industry avg: 28)\n"
        "Women in leadership roles: 43% (up from 38%)\n"
        "Internal promotions: 62% of open senior positions"
    )
    for para in body_ph6.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Book Antiqua"
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x33, 0x69, 0x1E)

    # Extra text box
    add_text_box(
        slide6,
        Inches(0.5), Inches(5.8), Inches(12), Inches(1.0),
        "\"Our people are our greatest asset.\" — CEO, Dr. Maria Santos",
        "Comic Sans MS", 15, bold=True,
        color_rgb=RGBColor(0x88, 0x00, 0x88),
        alignment=PP_ALIGN.CENTER
    )

    # ---- Slide 7: Innovation & R&D ----
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Innovation & R&D"
    slide7.shapes.title.text_frame.paragraphs[0].runs[0].font.name = "Segoe UI"
    slide7.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(34)
    slide7.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x96, 0x88)

    body_ph7 = slide7.placeholders[1]
    body_ph7.text = (
        "R&D investment: $340M (14.2% of revenue)\n"
        "Patents filed in 2024: 78 (approved: 52)\n"
        "Active research programs: 23 across 5 tech domains\n"
        "Partnerships with 12 universities and research institutes\n"
        "AI and ML integration embedded in 9 of 12 core products"
    )
    for para in body_ph7.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Consolas"
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x21, 0x21, 0x21)

    # ---- Slide 8: Sustainability ----
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Sustainability & ESG"
    slide8.shapes.title.text_frame.paragraphs[0].runs[0].font.name = "Rockwell"
    slide8.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(33)
    slide8.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x60, 0x0D)

    body_ph8 = slide8.placeholders[1]
    body_ph8.text = (
        "Carbon emissions reduced by 22% vs 2022 baseline\n"
        "100% renewable energy in all owned facilities\n"
        "Waste diversion rate: 91% (zero-landfill goal by 2026)\n"
        "Community investment: $18M across 42 programs\n"
        "ESG rating: AA (MSCI), 78/100 (Sustainalytics)"
    )
    for para in body_ph8.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Lucida Sans"
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0xE6, 0x5C, 0x00)

    add_text_box(
        slide8,
        Inches(8.5), Inches(3.8), Inches(4.0), Inches(2.0),
        "CDP Climate Score: A\nDow Jones Sustainability Index member\nFTSE4Good constituent",
        "Courier New", 13, bold=False,
        color_rgb=RGBColor(0x1B, 0x5E, 0x20),
        alignment=PP_ALIGN.LEFT
    )

    # ---- Slide 9: Outlook & Closing ----
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "2025 Outlook & Closing Remarks"
    slide9.shapes.title.text_frame.paragraphs[0].runs[0].font.name = "Arial Narrow"
    slide9.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(30)
    slide9.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xB7, 0x1C, 0x1C)

    body_ph9 = slide9.placeholders[1]
    body_ph9.text = (
        "Revenue guidance 2025: $2.65–$2.75 billion (+10-15%)\n"
        "Target operating margin: 20%+\n"
        "Strategic acquisitions pipeline: 3 targets under review\n"
        "New product launches planned: 6\n"
        "Thank you to all stakeholders for continued trust and partnership."
    )
    for para in body_ph9.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Tahoma"
            run.font.size = Pt(19)
            run.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    add_text_box(
        slide9,
        Inches(0.5), Inches(6.0), Inches(12), Inches(1.0),
        "Acme Corporation — Building Tomorrow, Today | www.acmecorp.example.com",
        "Impact", 13, bold=False,
        color_rgb=RGBColor(0x55, 0x55, 0x55),
        alignment=PP_ALIGN.CENTER
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
