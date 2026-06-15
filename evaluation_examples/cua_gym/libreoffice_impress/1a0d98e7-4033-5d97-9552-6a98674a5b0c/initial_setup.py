"""
Initial Setup: Create a 9-slide startup pitch deck with slide 4 titled 'Our Journey' but empty content.
Task ID: impress_ps_003
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_003'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

# Accent color for the presentation theme
ACCENT_COLOR = RGBColor(0x2E, 0x75, 0xB6)  # Professional blue
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
SUBTITLE_COLOR = RGBColor(0x66, 0x66, 0x66)


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_textbox(slide, left, top, width, height, text, font_size=28, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.LEFT):
    """Add a styled title text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_body_textbox(slide, left, top, width, height, lines, font_size=16, color=DARK_TEXT):
    """Add a multi-line body text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        if line:
            p.text = line
            run = p.runs[0]
            run.font.name = "Arial"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
        else:
            p.text = ""
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    # ── Slide 1: Title Slide ──
    slide1 = prs.slides.add_slide(blank_layout)
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    add_title_textbox(slide1, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
                      "NovaTech Solutions", font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_title_textbox(slide1, Inches(2), Inches(3.8), Inches(9), Inches(1.0),
                      "Revolutionizing Supply Chain Intelligence with AI", font_size=20, bold=False,
                      color=RGBColor(0xAA, 0xBB, 0xDD), alignment=PP_ALIGN.CENTER)
    add_title_textbox(slide1, Inches(3.5), Inches(5.2), Inches(6), Inches(0.6),
                      "Seed Round  |  Q2 2025  |  Confidential", font_size=14, bold=False,
                      color=RGBColor(0x88, 0x99, 0xBB), alignment=PP_ALIGN.CENTER)

    # ── Slide 2: The Problem ──
    slide2 = prs.slides.add_slide(blank_layout)
    add_title_textbox(slide2, Inches(0.8), Inches(0.5), Inches(10), Inches(1.0),
                      "The Problem", font_size=36, bold=True, color=ACCENT_COLOR)
    add_body_textbox(slide2, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5), [
        "• 73% of mid-market companies still rely on manual demand forecasting",
        "• Average $2.4M lost annually per company due to supply chain disruption",
        "• Existing tools require 6-12 month implementation cycles",
        "• No real-time visibility across multi-tier supplier networks",
        "• Legacy ERP systems lack predictive analytics capabilities",
    ], font_size=15, color=DARK_TEXT)
    add_body_textbox(slide2, Inches(7.0), Inches(1.8), Inches(5.5), Inches(3.0), [
        '"We spent $1.2M on SAP modules and still couldn\'t predict',
        'a chip shortage that cost us $8M in delayed shipments."',
        "",
        "— VP Operations, AutoPartsCo (Series B, 340 employees)",
    ], font_size=14, color=SUBTITLE_COLOR)

    # ── Slide 3: Our Solution ──
    slide3 = prs.slides.add_slide(blank_layout)
    add_title_textbox(slide3, Inches(0.8), Inches(0.5), Inches(10), Inches(1.0),
                      "Our Solution", font_size=36, bold=True, color=ACCENT_COLOR)
    add_body_textbox(slide3, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5), [
        "NovaTech delivers an AI-powered supply chain visibility platform:",
        "",
        "• Real-time risk scoring across 5,000+ global supplier nodes",
        "• Demand forecasting with 94% accuracy (vs. industry avg of 62%)",
        "• Plug-and-play integration: live in 2 weeks, not 6 months",
        "• Automated purchase order optimization reducing waste by 31%",
        "• Natural language query interface for non-technical supply chain managers",
    ], font_size=15, color=DARK_TEXT)

    # ── Slide 4: Our Journey (EMPTY content) ──
    slide4 = prs.slides.add_slide(blank_layout)
    add_title_textbox(slide4, Inches(0.8), Inches(0.5), Inches(10), Inches(1.0),
                      "Our Journey", font_size=36, bold=True, color=ACCENT_COLOR)
    # Content area is intentionally left EMPTY — no shapes, no timeline

    # ── Slide 5: Market Opportunity ──
    slide5 = prs.slides.add_slide(blank_layout)
    add_title_textbox(slide5, Inches(0.8), Inches(0.5), Inches(10), Inches(1.0),
                      "Market Opportunity", font_size=36, bold=True, color=ACCENT_COLOR)
    add_body_textbox(slide5, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5), [
        "Total Addressable Market (TAM): $47.2B",
        "Serviceable Addressable Market (SAM): $8.6B",
        "Serviceable Obtainable Market (SOM): $430M",
        "",
        "Key segments:",
        "• Mid-market manufacturing (100-2,000 employees)",
        "• Food & beverage distribution",
        "• Automotive parts suppliers",
        "• Pharmaceutical logistics",
    ], font_size=15, color=DARK_TEXT)
    add_body_textbox(slide5, Inches(7.0), Inches(1.8), Inches(5.5), Inches(3.0), [
        "Market growth: 18.4% CAGR through 2030",
        "Driven by:",
        "• Post-pandemic supply chain digitization",
        "• Regulatory compliance requirements",
        "• AI/ML cost reduction enabling mid-market adoption",
    ], font_size=14, color=SUBTITLE_COLOR)

    # ── Slide 6: Business Model ──
    slide6 = prs.slides.add_slide(blank_layout)
    add_title_textbox(slide6, Inches(0.8), Inches(0.5), Inches(10), Inches(1.0),
                      "Business Model", font_size=36, bold=True, color=ACCENT_COLOR)
    add_body_textbox(slide6, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5), [
        "SaaS subscription with usage-based pricing:",
        "",
        "• Starter: $2,500/mo — up to 500 supplier nodes, 3 users",
        "• Growth: $7,500/mo — up to 2,000 nodes, 10 users, API access",
        "• Enterprise: $25,000/mo — unlimited nodes, custom integrations, SLA",
        "",
        "Current metrics:",
        "  ARR: $1.8M  |  MRR growth: 22% MoM  |  Net retention: 134%",
        "  CAC payback: 8 months  |  Gross margin: 78%",
    ], font_size=15, color=DARK_TEXT)

    # ── Slide 7: Team ──
    slide7 = prs.slides.add_slide(blank_layout)
    add_title_textbox(slide7, Inches(0.8), Inches(0.5), Inches(10), Inches(1.0),
                      "Our Team", font_size=36, bold=True, color=ACCENT_COLOR)
    add_body_textbox(slide7, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5), [
        "Elena Rodriguez — CEO & Co-founder",
        "  Ex-McKinsey supply chain practice, MIT MBA",
        "",
        "James Okafor — CTO & Co-founder",
        "  Previously ML lead at Flexport, Stanford CS PhD",
        "",
        "Priya Sharma — VP Engineering",
        "  10 years at Oracle SCM, built teams of 50+",
        "",
        "Daniel Kim — VP Sales",
        "  Scaled revenue 0→$12M at ChainSight (acquired by SAP)",
    ], font_size=14, color=DARK_TEXT)
    add_body_textbox(slide7, Inches(7.0), Inches(1.8), Inches(5.5), Inches(3.0), [
        "Advisory Board:",
        "• Dr. Sarah Chen — Stanford AI Lab",
        "• Marco Benedetti — ex-CTO, Maersk Digital",
        "• Lisa Huang — Partner, Sequoia Capital",
    ], font_size=14, color=SUBTITLE_COLOR)

    # ── Slide 8: Financial Projections ──
    slide8 = prs.slides.add_slide(blank_layout)
    add_title_textbox(slide8, Inches(0.8), Inches(0.5), Inches(10), Inches(1.0),
                      "Financial Projections", font_size=36, bold=True, color=ACCENT_COLOR)
    # Simple table of projections
    table_shape = slide8.shapes.add_table(5, 4, Inches(1.0), Inches(1.8), Inches(10), Inches(3.0))
    table = table_shape.table
    headers = ["Metric", "2025 (Est.)", "2026 (Proj.)", "2027 (Proj.)"]
    data_rows = [
        ["ARR", "$3.6M", "$12.4M", "$38.0M"],
        ["Customers", "48", "165", "480"],
        ["Team Size", "22", "55", "120"],
        ["Burn Rate", "$180K/mo", "$350K/mo", "$520K/mo"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.name = "Arial"
    for r, row_data in enumerate(data_rows, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(13)
                run.font.name = "Arial"

    # ── Slide 9: Thank You / Contact ──
    slide9 = prs.slides.add_slide(blank_layout)
    fill = slide9.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    add_title_textbox(slide9, Inches(1.5), Inches(2.2), Inches(10), Inches(1.5),
                      "Thank You", font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_body_textbox(slide9, Inches(2), Inches(4.0), Inches(9), Inches(2.5), [
        "Elena Rodriguez  |  elena@novatech.ai  |  +1 (415) 555-0147",
        "www.novatech.ai",
        "",
        "NovaTech Solutions — San Francisco, CA",
    ], font_size=16, color=RGBColor(0xAA, 0xBB, 0xDD))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
