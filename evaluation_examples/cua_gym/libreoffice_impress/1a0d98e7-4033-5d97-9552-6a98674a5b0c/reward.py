"""
Reward Script: Milestone timeline on slide 4 of Startup_Pitch.pptx
Task ID: impress_ps_003
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): Horizontal right arrow shape on slide 4
  Component 2 (0.25): Five oval/circle marker shapes on slide 4
  Component 3 (0.35): All 5 milestone text labels present with correct content
  Component 4 (0.10): Milestone text font size is 12pt
  Component 5 (0.10): Ovals are evenly spaced horizontally along the arrow
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_003'

# The 5 milestones expected on the timeline
MILESTONES = [
    ("Founded", "Jan 2024"),
    ("MVP Launch", "Jun 2024"),
    ("First 1000 Users", "Dec 2024"),
    ("Series A", "Mar 2025"),
    ("Global Expansion", "Dec 2025"),
]


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"PRECONDITION FAIL: Need at least 4 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)

    # Collect shapes by type for analysis
    arrow_shapes = []
    oval_shapes = []
    text_shapes = []

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
                if shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW:
                    arrow_shapes.append(shape)
                elif shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.OVAL:
                    oval_shapes.append(shape)
            except Exception:
                pass
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            text_shapes.append(shape)

    # Component 1: Horizontal right arrow shape on slide 4 (0.20 points)
    try:
        if len(arrow_shapes) >= 1:
            arrow = arrow_shapes[0]
            # Verify it is horizontal (width > height)
            if arrow.width > arrow.height:
                print(f"PASS: Component 1 - Right arrow found, width={arrow.width}, height={arrow.height} (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 1 - Arrow found but not horizontal: width={arrow.width}, height={arrow.height}")
        else:
            print(f"FAIL: Component 1 - No RIGHT_ARROW shape found on slide 4. Found auto_shapes: {len(arrow_shapes)}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Five oval/circle marker shapes on slide 4 (0.25 points)
    try:
        num_ovals = len(oval_shapes)
        if num_ovals == 5:
            print(f"PASS: Component 2 - Found exactly 5 oval shapes (0.25 pts)")
            total_score += 0.25
        elif num_ovals >= 3:
            partial = 0.25 * (num_ovals / 5.0)
            print(f"PARTIAL: Component 2 - Found {num_ovals}/5 ovals ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 - Found {num_ovals} ovals, expected 5")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: All 5 milestone text labels present (0.35 points)
    try:
        # Gather all text from non-title text boxes on slide 4
        # Exclude the title shape ("Our Journey") by checking content
        milestone_texts = []
        for shape in text_shapes:
            if shape.has_text_frame:
                full_text = " ".join(p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip())
                if full_text and "Our Journey" not in full_text:
                    milestone_texts.append(full_text)

        matched = 0
        for label, date in MILESTONES:
            found = False
            for mt in milestone_texts:
                if label.lower() in mt.lower() and date.lower() in mt.lower():
                    found = True
                    break
            if found:
                matched += 1
                print(f"  FOUND milestone: '{label}' / '{date}'")
            else:
                print(f"  MISSING milestone: '{label}' / '{date}'")

        if matched == 5:
            print(f"PASS: Component 3 - All 5 milestones found (0.35 pts)")
            total_score += 0.35
        elif matched > 0:
            partial = 0.35 * (matched / 5.0)
            print(f"PARTIAL: Component 3 - {matched}/5 milestones found ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 - No milestones found among {len(milestone_texts)} text boxes")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Milestone text font size is 12pt (0.10 points)
    try:
        target_size = Pt(12)  # 152400 EMU
        size_correct_count = 0
        total_runs_checked = 0

        for shape in text_shapes:
            if shape.has_text_frame:
                full_text = " ".join(p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip())
                if "Our Journey" in full_text:
                    continue  # Skip title
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.text.strip():
                            total_runs_checked += 1
                            if r.font.size is not None and abs(r.font.size - target_size) < Pt(1):
                                size_correct_count += 1

        if total_runs_checked > 0 and size_correct_count == total_runs_checked:
            print(f"PASS: Component 4 - All {total_runs_checked} runs have ~12pt font (0.10 pts)")
            total_score += 0.10
        elif total_runs_checked > 0 and size_correct_count > 0:
            ratio = size_correct_count / total_runs_checked
            partial = 0.10 * ratio
            print(f"PARTIAL: Component 4 - {size_correct_count}/{total_runs_checked} runs at 12pt ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 - No milestone text runs at 12pt (checked {total_runs_checked} runs)")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Ovals are evenly spaced horizontally (0.10 points)
    try:
        if len(oval_shapes) == 5:
            # Sort ovals by left position
            sorted_ovals = sorted(oval_shapes, key=lambda s: s.left)
            centers = [s.left + s.width // 2 for s in sorted_ovals]

            # Calculate gaps between consecutive centers
            gaps = [centers[i+1] - centers[i] for i in range(len(centers) - 1)]

            if len(gaps) >= 2:
                avg_gap = sum(gaps) / len(gaps)
                # Check if all gaps are within 15% of average
                all_even = all(abs(g - avg_gap) / avg_gap < 0.15 for g in gaps)
                if all_even:
                    print(f"PASS: Component 5 - Ovals evenly spaced, gaps={[round(g) for g in gaps]} (0.10 pts)")
                    total_score += 0.10
                else:
                    print(f"FAIL: Component 5 - Ovals unevenly spaced, gaps={[round(g) for g in gaps]}, avg={round(avg_gap)}")
            else:
                print(f"FAIL: Component 5 - Not enough gaps to verify spacing")
        else:
            print(f"FAIL: Component 5 - Need 5 ovals for spacing check, found {len(oval_shapes)}")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
