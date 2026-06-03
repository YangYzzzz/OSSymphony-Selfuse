"""
Initial Setup: Add arrow shapes to slide 2 of ProcessFlow presentation
Task ID: impress_ndo_038
Domain: libreoffice_impress

Creates a multi-slide presentation with slide 2 titled 'Our Process'
but containing NO arrow shapes (pre-task state).
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_038'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "ProcessFlow"
    slide1.placeholders[1].text = "Q2 2025 Operations Review"

    # --- Slide 2: Title Only - "Our Process" (NO shapes - pre-task state) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide2.shapes.add_textbox(Cm(2), Cm(1), Cm(20), Cm(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Our Process"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # --- Slide 3: Team Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox3 = slide3.shapes.add_textbox(Cm(2), Cm(1), Cm(20), Cm(2))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Team Overview"
    p3.alignment = PP_ALIGN.LEFT
    run3 = p3.runs[0]
    run3.font.size = Pt(28)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # Add team member list
    members = [
        ("Sarah Chen", "Lead Engineer"),
        ("Marcus Johnson", "Product Manager"),
        ("Elena Rodriguez", "UX Designer"),
        ("David Kim", "Data Analyst"),
        ("Priya Patel", "QA Lead"),
    ]
    content_box = slide3.shapes.add_textbox(Cm(3), Cm(4), Cm(18), Cm(12))
    ctf = content_box.text_frame
    ctf.word_wrap = True
    for i, (name, role) in enumerate(members):
        if i == 0:
            p = ctf.paragraphs[0]
        else:
            p = ctf.add_paragraph()
        p.text = f"{name} - {role}"
        p.space_after = Pt(8)
        for r in p.runs:
            r.font.size = Pt(16)

    # --- Slide 4: Timeline ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox4 = slide4.shapes.add_textbox(Cm(2), Cm(1), Cm(20), Cm(2))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Project Timeline"
    p4.alignment = PP_ALIGN.LEFT
    run4 = p4.runs[0]
    run4.font.size = Pt(28)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    milestones = [
        "Phase 1: Requirements Gathering - March 2025",
        "Phase 2: Design & Prototyping - April 2025",
        "Phase 3: Development Sprint - May-June 2025",
        "Phase 4: Testing & QA - July 2025",
        "Phase 5: Launch & Deployment - August 2025",
    ]
    ms_box = slide4.shapes.add_textbox(Cm(3), Cm(4), Cm(18), Cm(12))
    ms_tf = ms_box.text_frame
    ms_tf.word_wrap = True
    for i, milestone in enumerate(milestones):
        if i == 0:
            p = ms_tf.paragraphs[0]
        else:
            p = ms_tf.add_paragraph()
        p.text = milestone
        p.space_after = Pt(10)
        for r in p.runs:
            r.font.size = Pt(14)

    # --- Slide 5: Summary ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox5 = slide5.shapes.add_textbox(Cm(2), Cm(1), Cm(20), Cm(2))
    tf5 = txBox5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Key Takeaways"
    p5.alignment = PP_ALIGN.LEFT
    run5 = p5.runs[0]
    run5.font.size = Pt(28)
    run5.font.bold = True
    run5.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    summary_box = slide5.shapes.add_textbox(Cm(3), Cm(4), Cm(18), Cm(10))
    stf = summary_box.text_frame
    stf.word_wrap = True
    points = [
        "Streamlined operations reduced cycle time by 23%",
        "Cross-functional collaboration improved with new workflow",
        "Customer satisfaction scores increased from 4.1 to 4.7",
        "Budget utilization at 94% - within target range",
    ]
    for i, point in enumerate(points):
        if i == 0:
            p = stf.paragraphs[0]
        else:
            p = stf.add_paragraph()
        p.text = point
        p.space_after = Pt(10)
        for r in p.runs:
            r.font.size = Pt(16)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
