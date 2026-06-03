"""
Reward Script: Add three right-pointing arrow shapes on slide 2 with specific positions, sizes, and colors.
Task ID: impress_ndo_038
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Three arrow auto-shapes exist on slide 2 with RIGHT_ARROW type
  Component 2 (0.25): Arrows positioned at correct Y coordinates (4cm, 8cm, 12cm)
  Component 3 (0.25): Arrows have correct dimensions (5cm wide x 2cm tall)
  Component 4 (0.25): Arrows have correct fill colors (#E74C3C, #F39C12, #27AE60)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_038'
FILE_NAME = f'{TASK_ID}.pptx'

# 1 cm = 360000 EMU
CM_TO_EMU = 360000
TOLERANCE = 0.02  # 2% relative tolerance for position/size checks

EXPECTED_ARROWS = [
    {'top_cm': 4.0, 'color': 'E74C3C'},
    {'top_cm': 8.0, 'color': 'F39C12'},
    {'top_cm': 12.0, 'color': '27AE60'},
]
EXPECTED_WIDTH_CM = 5.0
EXPECTED_HEIGHT_CM = 2.0


def is_approx(val_emu, expected_cm, tol=TOLERANCE):
    """Check if EMU value is approximately equal to expected cm value."""
    expected_emu = expected_cm * CM_TO_EMU
    if expected_emu == 0:
        return val_emu == 0
    return abs(val_emu - expected_emu) / expected_emu <= tol


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file has at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 2")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]  # 0-indexed

    # Collect all RIGHT_ARROW auto-shapes on slide 2
    arrow_shapes = []
    for shape in slide2.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # Check if it's a right arrow (auto_shape_type value 33 = RIGHT_ARROW)
                if shape.auto_shape_type is not None and shape.auto_shape_type == 33:
                    arrow_shapes.append(shape)
        except Exception:
            pass

    print(f"INFO: Found {len(arrow_shapes)} right-arrow shapes on slide 2")

    # Component 1: Three RIGHT_ARROW shapes exist on slide 2 (0.25 points)
    try:
        if len(arrow_shapes) == 3:
            print(f"PASS: Component 1 -- Exactly 3 right-arrow shapes on slide 2 (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 -- Expected 3 right-arrow shapes, found {len(arrow_shapes)}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    if len(arrow_shapes) < 3:
        # Cannot verify remaining components without 3 arrows
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Sort arrows by top position (Y) to match expected order
    arrow_shapes.sort(key=lambda s: s.top)

    # Component 2: Correct Y positions (4cm, 8cm, 12cm) (0.25 points)
    try:
        position_matches = 0
        for idx, (arrow, expected) in enumerate(zip(arrow_shapes, EXPECTED_ARROWS)):
            expected_top_cm = expected['top_cm']
            if is_approx(arrow.top, expected_top_cm):
                print(f"  Arrow {idx+1}: Y position OK (top={arrow.top/CM_TO_EMU:.2f}cm, expected={expected_top_cm}cm)")
                position_matches += 1
            else:
                print(f"  Arrow {idx+1}: Y position WRONG (top={arrow.top/CM_TO_EMU:.2f}cm, expected={expected_top_cm}cm)")

        if position_matches == 3:
            print(f"PASS: Component 2 -- All 3 arrows at correct Y positions (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 -- Only {position_matches}/3 arrows at correct Y positions")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Correct dimensions (5cm x 2cm) (0.25 points)
    try:
        size_matches = 0
        for idx, arrow in enumerate(arrow_shapes):
            w_ok = is_approx(arrow.width, EXPECTED_WIDTH_CM)
            h_ok = is_approx(arrow.height, EXPECTED_HEIGHT_CM)
            if w_ok and h_ok:
                print(f"  Arrow {idx+1}: Size OK ({arrow.width/CM_TO_EMU:.2f}cm x {arrow.height/CM_TO_EMU:.2f}cm)")
                size_matches += 1
            else:
                print(f"  Arrow {idx+1}: Size WRONG ({arrow.width/CM_TO_EMU:.2f}cm x {arrow.height/CM_TO_EMU:.2f}cm, expected 5.00cm x 2.00cm)")

        if size_matches == 3:
            print(f"PASS: Component 3 -- All 3 arrows have correct dimensions (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 -- Only {size_matches}/3 arrows have correct dimensions")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Correct fill colors (#E74C3C, #F39C12, #27AE60) (0.25 points)
    try:
        color_matches = 0
        for idx, (arrow, expected) in enumerate(zip(arrow_shapes, EXPECTED_ARROWS)):
            expected_color = expected['color'].upper()
            try:
                fill = arrow.fill
                if fill.type is not None and fill.type == 1:  # SOLID fill
                    actual_color = str(fill.fore_color.rgb).upper()
                    if actual_color == expected_color:
                        print(f"  Arrow {idx+1}: Color OK (#{actual_color})")
                        color_matches += 1
                    else:
                        print(f"  Arrow {idx+1}: Color WRONG (#{actual_color}, expected #{expected_color})")
                else:
                    print(f"  Arrow {idx+1}: Fill is not solid (type={fill.type}), expected solid fill #{expected_color}")
            except Exception as e:
                print(f"  Arrow {idx+1}: Could not read fill color: {e}")

        if color_matches == 3:
            print(f"PASS: Component 4 -- All 3 arrows have correct fill colors (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 4 -- Only {color_matches}/3 arrows have correct fill colors")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{FILE_NAME}'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
