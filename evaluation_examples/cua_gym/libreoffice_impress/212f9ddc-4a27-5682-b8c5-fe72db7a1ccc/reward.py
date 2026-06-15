"""
Reward Script: Build a fitness challenge presentation (8 slides)
Task ID: impress_wf_084
Domain: libreoffice_impress
Scoring:
  Component 1 (0.15): 8 slides total
  Component 2 (0.10): Slide 1 title contains '30-Day Fitness Challenge'
  Component 3 (0.15): Slide 2 has >= 4 card shapes (weekly themes)
  Component 4 (0.20): Slides 3-6 each have a table with >= 7 rows and 4 columns
  Component 5 (0.10): Slide 7 has a grid table (>= 5x5 for 30-day tracking)
  Component 6 (0.10): Slide 8 has >= 5 tip card shapes
  Component 7 (0.10): Push Down transitions on slides 3-6
  Component 8 (0.10): Color scheme uses #424242 background and white text
"""

import os
import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_084'
FILE_NAME = 'Fitness_Challenge.pptx'


def persist_app_state(domain):
    """Save any unsaved LibreOffice edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print(f"PERSIST: ctrl+s sent for {domain}")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def check_push_transition(pptx_path, slide_idx):
    """Check if a slide has a Push Down transition. slide_idx is 0-based."""
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            fname = f'ppt/slides/slide{slide_idx + 1}.xml'
            with zf.open(fname) as f:
                root = ET.parse(f).getroot()
                tr = root.find('.//p:transition', ns)
                if tr is not None:
                    push = tr.find('.//p:push', ns)
                    if push is not None:
                        direction = push.get('dir', '')
                        # Push Down: dir='d' or no dir (default is down)
                        if direction in ('d', ''):
                            return True
                return False
    except Exception:
        return False


def get_all_text_in_slide(slide):
    """Get all text content from a slide including nested shapes."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                txt = para.text.strip()
                if txt:
                    texts.append(txt)
    return texts


def count_auto_shapes(slide):
    """Count AUTO_SHAPE (rounded rectangles, etc.) shapes on a slide."""
    count = 0
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            count += 1
    return count


def get_tables(slide):
    """Get all table shapes from a slide."""
    tables = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            tables.append(shape.table)
    return tables


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Presentation has 8 slides (0.15 points)
    try:
        if num_slides == 8:
            print(f"PASS: Component 1 — 8 slides found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — Expected 8 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Need at least 8 slides for the remaining checks
    if num_slides < 8:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Slide 1 title contains '30-Day Fitness Challenge' (0.10 points)
    try:
        slide1_texts = get_all_text_in_slide(prs.slides[0])
        all_text = " ".join(slide1_texts).lower()
        if "30-day fitness challenge" in all_text:
            print(f"PASS: Component 2 — Slide 1 has '30-Day Fitness Challenge' (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 — Slide 1 text: {slide1_texts}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 2 has >= 4 card shapes (weekly theme cards) (0.15 points)
    try:
        card_count = count_auto_shapes(prs.slides[1])
        if card_count >= 4:
            print(f"PASS: Component 3 — Slide 2 has {card_count} card shapes (>= 4) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 3 — Slide 2 has {card_count} card shapes, need >= 4")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slides 3-6 each have exercise table with >= 7 rows x 4 cols (0.20 points)
    # Each valid table earns 0.05 points
    try:
        valid_tables = 0
        for slide_idx in range(2, 6):  # slides 3-6 (0-indexed: 2-5)
            tables = get_tables(prs.slides[slide_idx])
            if tables:
                t = tables[0]
                rows = len(t.rows)
                cols = len(t.columns)
                if rows >= 7 and cols >= 4:
                    valid_tables += 1
                    print(f"  Slide {slide_idx+1}: Table {rows}x{cols} — valid")
                else:
                    print(f"  Slide {slide_idx+1}: Table {rows}x{cols} — insufficient (need >=7 rows, >=4 cols)")
            else:
                print(f"  Slide {slide_idx+1}: No table found")

        table_score = valid_tables * 0.05
        if valid_tables == 4:
            print(f"PASS: Component 4 — All 4 week slides have valid exercise tables (0.20 pts)")
            total_score += 0.20
        elif valid_tables > 0:
            print(f"PARTIAL: Component 4 — {valid_tables}/4 valid tables ({table_score:.2f} pts)")
            total_score += table_score
        else:
            print(f"FAIL: Component 4 — No valid exercise tables found on slides 3-6")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slide 7 has a grid table for progress tracking (>= 5x5) (0.10 points)
    try:
        tables = get_tables(prs.slides[6])  # slide 7 (0-indexed: 6)
        if tables:
            t = tables[0]
            rows = len(t.rows)
            cols = len(t.columns)
            total_cells = rows * cols
            if total_cells >= 25 and rows >= 5 and cols >= 5:
                print(f"PASS: Component 5 — Slide 7 has {rows}x{cols} grid ({total_cells} cells) (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 5 — Slide 7 grid is {rows}x{cols} ({total_cells} cells), need >= 5x5")
        else:
            print(f"FAIL: Component 5 — No table found on slide 7")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Slide 8 has >= 5 tip card shapes (0.10 points)
    try:
        tip_count = count_auto_shapes(prs.slides[7])  # slide 8 (0-indexed: 7)
        if tip_count >= 5:
            print(f"PASS: Component 6 — Slide 8 has {tip_count} tip cards (>= 5) (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 6 — Slide 8 has {tip_count} tip cards, need >= 5")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Push Down transitions on slides 3-6 (0.10 points)
    # Each slide with correct transition earns 0.025 points
    try:
        transition_count = 0
        for slide_idx in range(2, 6):  # slides 3-6 (0-indexed: 2-5)
            if check_push_transition(file_path, slide_idx):
                transition_count += 1
                print(f"  Slide {slide_idx+1}: Push Down transition — present")
            else:
                print(f"  Slide {slide_idx+1}: Push Down transition — missing")

        if transition_count == 4:
            print(f"PASS: Component 7 — All 4 week slides have Push Down transitions (0.10 pts)")
            total_score += 0.10
        elif transition_count > 0:
            partial = transition_count * 0.025
            print(f"PARTIAL: Component 7 — {transition_count}/4 transitions ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 7 — No Push Down transitions found on slides 3-6")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    # Component 8: Color scheme - #424242 background and white text (0.10 points)
    try:
        bg_correct = 0
        white_text_found = False

        for slide_idx in range(len(prs.slides)):
            slide = prs.slides[slide_idx]
            fill = slide.background.fill
            if fill.type == 1:  # SOLID
                try:
                    bg_rgb = str(fill.fore_color.rgb).upper()
                    if bg_rgb == "424242":
                        bg_correct += 1
                except:
                    pass

            # Check for white text on this slide
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                try:
                                    if run.font.color.type is not None:
                                        rgb = str(run.font.color.rgb).upper()
                                        if rgb == "FFFFFF":
                                            white_text_found = True
                                except:
                                    pass

        # Need majority of slides with #424242 background AND white text somewhere
        if bg_correct >= 6 and white_text_found:
            print(f"PASS: Component 8 — {bg_correct}/8 slides have #424242 bg, white text found (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 8 — {bg_correct}/8 slides with #424242 bg, white text: {white_text_found}")
    except Exception as e:
        print(f"ERROR: Component 8 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{FILE_NAME}'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
