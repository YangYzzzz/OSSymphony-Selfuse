"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 157 is the odd one out—its heading is stuck in Liberation Sans 32 pt. In LibreOffice Impress, how do I switch that single title to Calibri 40 pt, and if Calibri isn’t installed on the machine, pick the nearest available font at the same 40 pt size?
Generated: 2025-09-10 16:26:25
Status: success
Model: azure-o3
Total Steps: 2
"""

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
import os
from collections import Counter

FILE_PATH = "/home/user/slide_157_is_the_odd_one_outits_heading_is_stuck_in_liberation_sans_32_pt_in_libreoffice_impress_how_golden.pptx"

def _get_title_shape(slide):
    """Return the title placeholder if it exists; otherwise the first text frame shape."""
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE and shape.has_text_frame:
            return shape
    for shape in slide.shapes:  # fallback
        if shape.has_text_frame:
            return shape
    return None

def verify_slide_157_title_font(file_path: str) -> float:
    """Verify that slide 157's title is 40 pt and uses Calibri (or any non-Liberation Sans font).

    Scoring (progressive):
      • 0.5 pts  — Title font size is exactly 40 pt (rounded).
      • 0.25 pts — Title font family is *not* Liberation Sans.
      • 0.25 pts — Title font family is Calibri (bonus on top of the previous 0.25).
      → Perfect completion = 1.0
    """
    print(f"Loading presentation: {file_path}")

    if not os.path.exists(file_path):
        print("✗ File does not exist")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error opening PPTX: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Slide 157 is index 156 (0-based)
    idx = 156
    if idx >= len(prs.slides):
        print(f"✗ Slide 157 not found (presentation only has {len(prs.slides)} slides)")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[idx]
    title_shape = _get_title_shape(slide)
    if not title_shape:
        print("✗ Could not locate a title/text shape on slide 157")
        print("REWARD: 0.0")
        return 0.0

    # Collect font names & sizes
    font_names, font_sizes = [], []
    for para in title_shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.name:
                font_names.append(run.font.name)
            if run.font.size:
                try:
                    font_sizes.append(run.font.size.pt)
                except Exception:
                    pass  # size might be EMU; ignore if conversion fails

    print(f"Collected font names: {font_names}")
    print(f"Collected font sizes (pt): {font_sizes}")

    score = 0.0

    # 1) Size check (0.5 pts)
    if font_sizes:
        most_common_size = Counter(round(s) for s in font_sizes).most_common(1)[0][0]
        print(f"Most common rounded size: {most_common_size} pt")
        if most_common_size == 40:
            score += 0.5
            print("✓ Title font size set to 40 pt (0.5 points)")
        else:
            print("✗ Title font size is not 40 pt (0 points)")
    else:
        print("✗ No explicit font sizes detected (0 points)")

    # 2) Font family checks (0.25 + 0.25 pts)
    if font_names:
        most_common_font = Counter(font_names).most_common(1)[0][0]
        font_lc = most_common_font.lower()
        print(f"Most common font name: {most_common_font}")

        if font_lc != "liberation sans":
            score += 0.25
            print("✓ Title font is not Liberation Sans (0.25 points)")
            if font_lc == "calibri":
                score += 0.25
                print("✓ Title font is Calibri (additional 0.25 points)")
            else:
                print("Title font is an alternative to Calibri (no extra Calibri points)")
        else:
            print("✗ Title font is Liberation Sans (0 points)")
    else:
        print("✗ No explicit font names detected (0 points)")

    final_score = min(score, 1.0)
    print(f"Total score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score

if __name__ == "__main__":
    verify_slide_157_title_font(FILE_PATH)

