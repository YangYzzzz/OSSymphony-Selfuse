"""
Reward Script: Photo album layout on slide 6 with 4 images
Task ID: impress_tm_093
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): 4 pictures present on slide 6
  Component 2 (0.25): Each picture is 6cm x 6cm
  Component 3 (0.25): Pictures in horizontal row with ~1cm gaps
  Component 4 (0.20): Each picture has thin black border (~0.5pt)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_093'

# Constants
CM_TO_EMU = 360000  # 1 cm = 360000 EMU
TARGET_SIZE = 6 * CM_TO_EMU  # 6cm = 2160000 EMU
TARGET_GAP = 1 * CM_TO_EMU   # 1cm = 360000 EMU
SIZE_TOLERANCE = 0.05         # 5% tolerance on size
GAP_TOLERANCE = 0.15          # 15% tolerance on gap
BORDER_WIDTH_TOLERANCE = 0.3  # 30% tolerance on border width (0.5pt = 6350 EMU)
TARGET_BORDER_WIDTH = 6350    # 0.5pt in EMU


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: File has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[5]  # slide 6 (0-indexed)

    # Collect all PICTURE shapes on slide 6
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]

    # Component 1: 4 pictures present on slide 6 (0.30 points)
    try:
        pic_count = len(pictures)
        if pic_count == 4:
            print(f"PASS: Component 1 -- 4 pictures found on slide 6 (0.30 pts)")
            total_score += 0.30
        elif pic_count >= 1:
            # Partial credit: at least some pictures inserted
            partial = 0.30 * (min(pic_count, 4) / 4)
            print(f"PARTIAL: Component 1 -- {pic_count} pictures found (expected 4), awarding {partial:.2f} pts")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- No pictures found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # If no pictures found, remaining components cannot pass
    if len(pictures) == 0:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: Each picture is 6cm x 6cm (0.25 points)
    try:
        correct_size_count = 0
        for pic in pictures:
            w_ok = abs(pic.width - TARGET_SIZE) / TARGET_SIZE <= SIZE_TOLERANCE
            h_ok = abs(pic.height - TARGET_SIZE) / TARGET_SIZE <= SIZE_TOLERANCE
            if w_ok and h_ok:
                correct_size_count += 1
            else:
                print(f"  INFO: {pic.name} size ({pic.width},{pic.height}) expected ~({TARGET_SIZE},{TARGET_SIZE})")

        if correct_size_count == len(pictures) and len(pictures) >= 4:
            print(f"PASS: Component 2 -- All 4 pictures are ~6cm x 6cm (0.25 pts)")
            total_score += 0.25
        elif correct_size_count > 0:
            partial = 0.25 * (correct_size_count / max(len(pictures), 4))
            print(f"PARTIAL: Component 2 -- {correct_size_count}/{len(pictures)} pictures correct size, awarding {partial:.2f} pts")
            total_score += partial
        else:
            print(f"FAIL: Component 2 -- No pictures have correct 6cm x 6cm size")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Pictures arranged in horizontal row with ~1cm gaps (0.25 points)
    try:
        if len(pictures) >= 2:
            # Sort pictures by left position
            sorted_pics = sorted(pictures, key=lambda p: p.left)

            # Check same vertical position (horizontal row)
            tops = [p.top for p in sorted_pics]
            same_row = all(abs(t - tops[0]) / max(abs(tops[0]), 1) <= 0.05 for t in tops)

            # Check gaps between consecutive pictures
            gaps_ok = 0
            total_gaps = len(sorted_pics) - 1
            for i in range(total_gaps):
                gap = sorted_pics[i + 1].left - (sorted_pics[i].left + sorted_pics[i].width)
                if abs(gap - TARGET_GAP) / TARGET_GAP <= GAP_TOLERANCE:
                    gaps_ok += 1
                else:
                    print(f"  INFO: Gap between pic {i} and {i+1}: {gap} EMU (expected ~{TARGET_GAP})")

            if same_row and gaps_ok == total_gaps and len(sorted_pics) >= 4:
                print(f"PASS: Component 3 -- Pictures in horizontal row with ~1cm gaps (0.25 pts)")
                total_score += 0.25
            elif same_row and gaps_ok > 0:
                partial = 0.25 * (0.5 + 0.5 * gaps_ok / max(total_gaps, 1))
                print(f"PARTIAL: Component 3 -- Row OK, {gaps_ok}/{total_gaps} gaps correct, awarding {partial:.2f} pts")
                total_score += partial
            elif same_row:
                print(f"PARTIAL: Component 3 -- Same row but gaps incorrect, awarding 0.10 pts")
                total_score += 0.10
            else:
                print(f"FAIL: Component 3 -- Pictures not in a horizontal row")
        else:
            print(f"FAIL: Component 3 -- Need at least 2 pictures to check arrangement")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Each picture has thin black border (~0.5pt) (0.20 points)
    try:
        border_ok_count = 0
        for pic in pictures:
            has_border = False
            try:
                line_width = pic.line.width
                if line_width is not None and line_width > 0:
                    # Check width is approximately 0.5pt (6350 EMU) — accept thin borders up to ~2pt
                    width_ok = line_width <= 25400  # up to 2pt is "thin"
                    # Check color is black
                    color_ok = False
                    try:
                        if pic.line.color.type is not None:
                            rgb = str(pic.line.color.rgb)
                            color_ok = rgb == '000000'
                        else:
                            color_ok = False
                    except Exception:
                        color_ok = False

                    if width_ok and color_ok:
                        has_border = True
                    else:
                        if not width_ok:
                            print(f"  INFO: {pic.name} border width {line_width} EMU not thin")
                        if not color_ok:
                            print(f"  INFO: {pic.name} border color not black")
            except Exception:
                pass

            if has_border:
                border_ok_count += 1

        if border_ok_count == len(pictures) and len(pictures) >= 4:
            print(f"PASS: Component 4 -- All pictures have thin black border (0.20 pts)")
            total_score += 0.20
        elif border_ok_count > 0:
            partial = 0.20 * (border_ok_count / max(len(pictures), 4))
            print(f"PARTIAL: Component 4 -- {border_ok_count}/{len(pictures)} pictures have correct border, awarding {partial:.2f} pts")
            total_score += partial
        else:
            print(f"FAIL: Component 4 -- No pictures have thin black border")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
persist_app_state()
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
