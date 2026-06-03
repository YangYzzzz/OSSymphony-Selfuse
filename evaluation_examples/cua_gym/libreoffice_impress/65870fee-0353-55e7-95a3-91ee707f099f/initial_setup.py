"""
Initial Setup: Create event recap presentation with 8 slides, slide 6 titled 'Event Highlights'
Task ID: impress_tm_093
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_093'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
PHOTOS_DIR = f'{WORKDIR}/photos'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_sample_images():
    """Create 4 sample event photos."""
    os.makedirs(PHOTOS_DIR, exist_ok=True)
    # Create distinct colored images to simulate event photos
    colors_and_labels = [
        ((180, 60, 60), "Event Photo 1"),
        ((60, 130, 180), "Event Photo 2"),
        ((80, 160, 80), "Event Photo 3"),
        ((180, 140, 50), "Event Photo 4"),
    ]
    for i, (color, label) in enumerate(colors_and_labels, 1):
        img = Image.new('RGB', (600, 600), color)
        path = os.path.join(PHOTOS_DIR, f'event{i}.jpg')
        img.save(path, 'JPEG', quality=85)
        print(f'Created: {path}')


def create_initial():
    prs = Presentation()
    # Set widescreen 13.333 x 7.5 inches (16:9) so 4 x 6cm images + gaps fit
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Annual Event Recap 2025"
    slide1.placeholders[1].text = "Marketing Department Review"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Overview of Key Events"
    p = body2.add_paragraph()
    p.text = "Attendance and Engagement Metrics"
    p = body2.add_paragraph()
    p.text = "Budget Summary"
    p = body2.add_paragraph()
    p.text = "Event Highlights"
    p = body2.add_paragraph()
    p.text = "Lessons Learned"
    p = body2.add_paragraph()
    p.text = "Next Steps for 2026"

    # --- Slide 3: Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "2025 Event Overview"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Total Events Organized: 12"
    p = body3.add_paragraph()
    p.text = "Total Attendees: 4,850"
    p = body3.add_paragraph()
    p.text = "Customer Satisfaction Score: 4.7/5.0"
    p = body3.add_paragraph()
    p.text = "Revenue Generated: $1.2M"

    # --- Slide 4: Attendance Metrics ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Attendance & Engagement"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Q1: Spring Launch - 620 attendees"
    p = body4.add_paragraph()
    p.text = "Q2: Summer Conference - 1,200 attendees"
    p = body4.add_paragraph()
    p.text = "Q3: Tech Summit - 1,850 attendees"
    p = body4.add_paragraph()
    p.text = "Q4: Year-End Gala - 1,180 attendees"

    # --- Slide 5: Budget Summary ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Budget Summary"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Total Budget: $450,000"
    p = body5.add_paragraph()
    p.text = "Spent: $412,300 (91.6%)"
    p = body5.add_paragraph()
    p.text = "Venue Costs: $185,000"
    p = body5.add_paragraph()
    p.text = "Catering: $98,500"
    p = body5.add_paragraph()
    p.text = "Marketing & Promotion: $72,800"
    p = body5.add_paragraph()
    p.text = "Remaining: $37,700"

    # --- Slide 6: Event Highlights (ONLY TITLE, no images) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide6.shapes.add_textbox(Cm(2), Cm(1), Cm(30), Cm(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Event Highlights"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 7: Lessons Learned ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Lessons Learned"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "Earlier venue booking saves 15-20% on costs"
    p = body7.add_paragraph()
    p.text = "Hybrid format increased reach by 40%"
    p = body7.add_paragraph()
    p.text = "Post-event surveys within 48 hours get 3x response rate"
    p = body7.add_paragraph()
    p.text = "Social media promotion should start 6 weeks before event"

    # --- Slide 8: Next Steps ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Next Steps for 2026"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Increase event count to 15"
    p = body8.add_paragraph()
    p.text = "Launch dedicated event mobile app"
    p = body8.add_paragraph()
    p.text = "Expand international presence with 3 overseas events"
    p = body8.add_paragraph()
    p.text = "Partner with industry leaders for co-hosted summits"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


# Create photos first, then presentation
create_sample_images()
create_initial()
