"""
FINAL REWARD SCRIPT - SUCCESS
Task: Convert all footnotes to endnotes.
Generated: 2025-10-17 07:53:18
Status: success
Model: azure-o3
Total Steps: 7
"""

import re
from pptx import Presentation


def verify_footnotes_converted_to_endnotes(file_path: str) -> float:
    """Verify that all footnotes have been converted to endnotes.

    Scoring (progressive):
        0.4 – A dedicated slide titled exactly 'Endnotes' exists.
        0.2 – That Endnotes slide contains at least one enumerated (numbered) line.
        0.4 – No residual numbered footnotes remain on any other slide (heuristic: numbered
              lines located near the bottom 30% of a slide).
    Returns:
        float: score between 0.0 and 1.0
    """

    # Regular expression for lines that look like numbered footnotes (e.g., "1. text", "2) text")
    enum_pattern = re.compile(r"^\s*\d+\s*[\.)]\s*", re.UNICODE)

    # Attempt to load the presentation ---------------------------------------------------------
    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded presentation with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Could not open PPTX: {e}")
        print("REWARD: 0.0")
        return 0.0

    max_score = 1.0
    score = 0.0

    # 1. Locate the slide titled exactly 'Endnotes' --------------------------------------------
    endnotes_idx = None
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip().lower() == "endnotes":
                endnotes_idx = idx
                break
        if endnotes_idx is not None:
            break

    if endnotes_idx is None:
        print("✗ No slide titled 'Endnotes' found")
    else:
        print(f"✓ Found 'Endnotes' slide at position {endnotes_idx + 1}")
        score += 0.4

    # 2. Ensure the Endnotes slide actually contains numbered lines -----------------------------
    enumerated_lines = []
    if endnotes_idx is not None:
        end_slide = prs.slides[endnotes_idx]
        for shape in end_slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for line in shape.text.splitlines():
                if enum_pattern.match(line):
                    enumerated_lines.append(line.strip())
        if enumerated_lines:
            print(f"✓ Endnotes slide contains {len(enumerated_lines)} enumerated line(s):")
            for ln in enumerated_lines:
                print("   ", ln)
            score += 0.2
        else:
            print("✗ Endnotes slide has no enumerated footnote lines")

    # 3. Detect residual numbered footnotes on other slides ------------------------------------
    residual_footnotes = []
    slide_height = prs.slide_height
    for idx, slide in enumerate(prs.slides):
        if idx == endnotes_idx:
            continue  # skip the Endnotes slide itself
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            # Heuristic: footnotes are typically near the bottom (shape.top > 70% height)
            near_bottom = False
            try:
                near_bottom = shape.top > 0.7 * slide_height
            except Exception:
                pass  # if shape.top isn't available, ignore the positional heuristic
            for line in shape.text.splitlines():
                if enum_pattern.match(line):
                    residual_footnotes.append((idx + 1, line.strip(), near_bottom))

    if not residual_footnotes:
        print("✓ No residual numbered footnotes detected on other slides")
        score += 0.4
    else:
        bottom_count = sum(1 for _, _, nb in residual_footnotes if nb)
        print(f"✗ Detected {len(residual_footnotes)} potential footnote(s) remaining on content slides")
        for slide_no, txt, nb in residual_footnotes:
            loc = "(bottom)" if nb else "(not bottom)"
            print(f"   Slide {slide_no}: '{txt}' {loc}")
        # If no matches are in bottom area, grant partial credit (0.2)
        if bottom_count == 0:
            score += 0.2

    final_score = min(score, max_score)
    print(f"Total Score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score


if __name__ == "__main__":
    # Path provided in the task context
    FILE_PATH = "/home/user/convert_all_footnotes_to_endnotes.pptx"
    verify_footnotes_converted_to_endnotes(FILE_PATH)
