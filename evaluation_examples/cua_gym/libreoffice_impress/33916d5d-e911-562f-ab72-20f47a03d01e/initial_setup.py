"""
Initial Setup: Architectural portfolio with facade photo at ~55% slide area on slide 1
Task ID: osworld_impress_image_fill_slide_011
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import io
import struct
import zlib

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_image_fill_slide_011'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
IMG_PATH = f'{WORKDIR}/{TASK_ID}_facade.jpg'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_facade_image():
    """Create a realistic-looking facade photo using Pillow."""
    try:
        from PIL import Image, ImageDraw, ImageFilter
        import random

        random.seed(42)

        # Create a 1200x800 image representing a building facade
        width, height = 1200, 800
        img = Image.new('RGB', (width, height), color=(200, 190, 175))
        draw = ImageDraw.Draw(img)

        # Sky gradient (top portion)
        for y in range(200):
            r = int(135 + (y / 200) * 40)
            g = int(180 + (y / 200) * 20)
            b = int(235 - (y / 200) * 30)
            draw.line([(0, y), (width, y)], fill=(r, g, b))

        # Building body - warm sandstone color
        draw.rectangle([50, 180, 1150, 760], fill=(210, 195, 170))

        # Architectural facade details - horizontal bands
        draw.rectangle([50, 180, 1150, 210], fill=(190, 175, 150))
        draw.rectangle([50, 380, 1150, 400], fill=(190, 175, 150))
        draw.rectangle([50, 570, 1150, 590], fill=(190, 175, 150))

        # Windows - row 1
        window_color = (80, 100, 130)
        highlight_color = (150, 170, 200)
        for col in range(6):
            x = 100 + col * 175
            # Row 1 windows
            draw.rectangle([x, 230, x+100, 350], fill=window_color)
            draw.rectangle([x, 235, x+50, 290], fill=highlight_color)
            # Row 2 windows
            draw.rectangle([x, 420, x+100, 540], fill=window_color)
            draw.rectangle([x, 425, x+50, 480], fill=highlight_color)

        # Ground floor arched entryway
        draw.rectangle([450, 620, 750, 760], fill=(160, 145, 120))
        # Arch
        draw.ellipse([450, 590, 750, 680], fill=(160, 145, 120))
        draw.rectangle([480, 630, 720, 760], fill=(100, 85, 65))

        # Columns
        for x in [200, 400, 800, 1000]:
            draw.rectangle([x, 180, x+30, 760], fill=(225, 215, 195))

        # Cornice at top
        draw.rectangle([30, 165, 1170, 185], fill=(170, 155, 130))
        draw.rectangle([20, 155, 1180, 167], fill=(185, 170, 145))

        # Ground/pavement
        draw.rectangle([0, 760, width, height], fill=(140, 135, 125))
        for x in range(0, width, 80):
            draw.line([(x, 760), (x, height)], fill=(120, 115, 105), width=1)
        for y in range(760, height, 25):
            draw.line([(0, y), (width, y)], fill=(120, 115, 105), width=1)

        # Slight blur for realism
        img = img.filter(ImageFilter.SMOOTH)

        img.save(IMG_PATH, 'JPEG', quality=85)
        print(f'Facade image created: {IMG_PATH}')
        return True
    except ImportError:
        # Fallback: create minimal JPEG manually
        return create_minimal_jpeg()


def create_minimal_jpeg():
    """Create a minimal valid JPEG file as fallback."""
    # Minimal 100x67 JPEG with building-like colors (gray/tan)
    try:
        from PIL import Image
        img = Image.new('RGB', (600, 400), color=(200, 190, 170))
        img.save(IMG_PATH, 'JPEG', quality=85)
        print(f'Facade image created (minimal): {IMG_PATH}')
        return True
    except Exception as e:
        print(f'Warning: could not create image: {e}')
        return False


def create_initial():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # Create facade image first
    create_facade_image()

    prs = Presentation()

    # Standard widescreen 16:9 slide dimensions
    slide_width = prs.slide_width    # 9144000 EMU = 10 inches
    slide_height = prs.slide_height  # 6858000 EMU = 7.5 inches

    # ── Slide 1: Title slide with facade photo (~55% of slide area, centered) ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Add title text box
    txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Meridian Architecture Studio"
    run.font.name = "Calibri"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # Add subtitle
    txBox2 = slide1.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.6))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Portfolio 2024 — Urban Facades & Public Spaces"
    run2.font.name = "Calibri"
    run2.font.size = Pt(16)
    run2.font.italic = True
    run2.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)

    # Add facade photo centered, ~55% of slide area
    # Slide is 10" x 7.5" = 75 sq inches; 55% = ~41.25 sq inches
    # Keep 4:3 ratio for the image: ~7.4" x 5.6" ≈ 41.4 sq inches
    img_w = Inches(7.4)
    img_h = Inches(5.0)
    img_left = (slide_width - img_w) // 2
    img_top = Inches(1.9)

    if os.path.exists(IMG_PATH):
        pic = slide1.shapes.add_picture(IMG_PATH, img_left, img_top, img_w, img_h)
        print(f'Slide 1: image added at center, size {img_w} x {img_h} EMU')
    else:
        # Fallback: add a rectangle as placeholder
        shape = slide1.shapes.add_shape(1, img_left, img_top, img_w, img_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xAA, 0x99, 0x88)

    # ── Slide 2: Project Overview ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide2.shapes.title.text = "Project Overview"
    tf3 = slide2.placeholders[1].text_frame
    tf3.text = "Harborview Cultural Center — Phase 1 Complete"
    paras = [
        "Location: 14 Maritime Boulevard, Auckland",
        "Client: Auckland City Council",
        "Gross Floor Area: 8,450 m²",
        "Completion: March 2024",
        "Budget: NZD 24.6 Million",
        "Sustainability: 5-Star Green Star Certified",
    ]
    for text in paras:
        p = tf3.add_paragraph()
        p.text = text
        p.level = 1

    # ── Slide 3: Design Concept ──
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Design Concept"
    tf4 = slide3.placeholders[1].text_frame
    tf4.text = "Inspired by Auckland's volcanic landscape and maritime heritage"
    concepts = [
        "Basalt-textured precast concrete panels reference Rangitoto Island",
        "Generous public forecourt activates Maritime Boulevard frontage",
        "Sky-bridges connect rooftop terrace with adjacent civic plaza",
        "North-facing glazing maximises passive solar gain in winter months",
        "Double-skin facade reduces solar heat load by 40% in summer",
    ]
    for text in concepts:
        p = tf4.add_paragraph()
        p.text = text
        p.level = 1

    # ── Slide 4: Material Palette ──
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Material Palette"
    tf5 = slide4.placeholders[1].text_frame
    tf5.text = "Durable, locally-sourced materials for longevity and character"
    materials = [
        "Exterior: Textured precast concrete — Eurostone Basalt Mix",
        "Cladding: Zinc standing-seam panels — VM Zinc Quartz-Zinc",
        "Glazing: Triple-glazed unitised curtain wall — Schüco AWS 90.SI",
        "Interior: Exposed structural timber — Engineered LVL pine",
        "Flooring: Terrazzo with recycled glass aggregate",
        "Landscaping: Pohutukawa grove and native plantings",
    ]
    for text in materials:
        p = tf5.add_paragraph()
        p.text = text
        p.level = 1

    # ── Slide 5: Project Timeline ──
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Project Timeline"
    tf6 = slide5.placeholders[1].text_frame
    tf6.text = "Key Milestones — Harborview Cultural Center"
    timeline = [
        "Feb 2021 — Site survey and geotechnical investigation",
        "Sep 2021 — Resource consent approved by Auckland Council",
        "Jan 2022 — Detailed design and documentation complete",
        "Apr 2022 — Construction tender awarded to Fletcher Building",
        "Jun 2022 — Ground breaking and foundation works commence",
        "Oct 2022 — Structural steel frame topped out",
        "Mar 2023 — Facade installation and weathertight",
        "Nov 2023 — Interior fit-out and commissioning",
        "Mar 2024 — Practical completion and handover",
    ]
    for text in timeline:
        p = tf6.add_paragraph()
        p.text = text
        p.level = 1

    prs.save(OUTPUT)
    print(f'Initial presentation created: {OUTPUT}')
    print(f'Slides: {len(prs.slides)}')
    print(f'Slide 1 image dimensions: ~55% of slide area (centered)')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
