"""
Reward Script: Scale image on slide 1 to fill the full slide (full-bleed background)
Task ID: osworld_impress_image_fill_slide_011
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Image height fills full slide height (covers entire slide vertically)
  Component 2 (0.3): Image width covers full slide width (no uncovered horizontal area)
  Component 3 (0.3): Image aspect ratio is maintained (within 1% of native 1.5:1 ratio)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_image_fill_slide_011'

# Native image aspect ratio (1200x800 pixels)
NATIVE_ASPECT_RATIO = 1.5
TOLERANCE = 0.01  # 1% relative tolerance for dimension checks


def is_approx_equal(val1, val2, tolerance=TOLERANCE):
    """Check if two values are approximately equal within relative tolerance."""
    if val1 == val2:
        return True
    if max(abs(val1), abs(val2)) == 0:
        return True
    return abs(val1 - val2) / max(abs(val1), abs(val2)) <= tolerance


def verify_task(file_path):
    """
    Verify task completion: image on slide 1 scaled to full-bleed background.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: 5-slide presentation with correct slide dimensions
    slide_width = prs.slide_width   # 9144000 EMU = 10 inches
    slide_height = prs.slide_height  # 6858000 EMU = 7.5 inches

    if len(prs.slides) < 1:
        print("CRITICAL: Presentation has no slides")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]

    # Find the picture on slide 1
    picture_shape = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture_shape = shape
            break

    if picture_shape is None:
        print("FAIL: No picture found on slide 1")
        print("REWARD: 0.0")
        return 0.0

    img_left = picture_shape.left
    img_top = picture_shape.top
    img_width = picture_shape.width
    img_height = picture_shape.height

    print(f"INFO: Slide size: {slide_width} x {slide_height} EMU ({slide_width/914400:.3f} x {slide_height/914400:.3f} inches)")
    print(f"INFO: Image: left={img_left} top={img_top} width={img_width} height={img_height} EMU")
    print(f"INFO: Image: left={img_left/914400:.4f} top={img_top/914400:.4f} w={img_width/914400:.4f} h={img_height/914400:.4f} inches")

    # Component 1: Image height fills the full slide height (0.4 points)
    # The image must cover the full slide height — height should equal slide_height
    try:
        height_matches_slide = is_approx_equal(img_height, slide_height, tolerance=0.01)
        top_at_zero = is_approx_equal(img_top, 0, tolerance=0.02) or img_top <= 0

        if height_matches_slide and top_at_zero:
            print(f"PASS: Component 1 — Image height {img_height} matches slide height {slide_height} and top={img_top} (0.4 pts)")
            total_score += 0.4
        else:
            if not height_matches_slide:
                print(f"FAIL: Component 1 — Image height {img_height} does not match slide height {slide_height} (expected approx equal)")
            if not top_at_zero:
                print(f"FAIL: Component 1 — Image top {img_top} is not at 0 (should be top-aligned to fill slide)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Image width covers full slide width (no uncovered horizontal area) (0.3 points)
    # The image should be wide enough so that the entire slide width is covered
    # i.e., left <= 0 and left + width >= slide_width
    try:
        covers_left = img_left <= 0  # image starts at or before left edge
        right_edge = img_left + img_width
        covers_right = right_edge >= slide_width  # image extends to or past right edge

        if covers_left and covers_right:
            print(f"PASS: Component 2 — Image covers full slide width: left={img_left/914400:.4f}in, right={right_edge/914400:.4f}in (slide={slide_width/914400:.4f}in) (0.3 pts)")
            total_score += 0.3
        else:
            if not covers_left:
                print(f"FAIL: Component 2 — Image left edge {img_left/914400:.4f}in does not reach/exceed slide left edge (should be <= 0)")
            if not covers_right:
                print(f"FAIL: Component 2 — Image right edge {right_edge/914400:.4f}in does not reach slide right edge {slide_width/914400:.4f}in")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Image aspect ratio is maintained (within 1% of native 1.5:1 ratio) (0.3 points)
    # The task requires "maintaining aspect ratio" — check that image width/height ~ 1.5
    try:
        if img_height > 0:
            actual_ratio = img_width / img_height
            ratio_ok = is_approx_equal(actual_ratio, NATIVE_ASPECT_RATIO, tolerance=0.01)
            if ratio_ok:
                print(f"PASS: Component 3 — Image aspect ratio {actual_ratio:.4f} matches native ratio {NATIVE_ASPECT_RATIO} (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 3 — Image aspect ratio {actual_ratio:.4f} does not match native ratio {NATIVE_ASPECT_RATIO} (tolerance 1%)")
        else:
            print("FAIL: Component 3 — Image height is zero, cannot compute aspect ratio")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Run verification against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
