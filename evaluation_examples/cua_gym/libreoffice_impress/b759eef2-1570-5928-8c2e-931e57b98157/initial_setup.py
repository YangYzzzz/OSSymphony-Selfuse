"""
Initial Setup: Create academic lab results presentation without watermark
Task ID: impress_cross_acad_028
Domain: libreoffice_impress
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_cross_acad_028'
# The actual file the agent will work on (as referenced in the task instruction)
TASK_FILE = '/home/user/Documents/lab_results_talk.odp'
# The initial file stored in workdir for tracking
OUTPUT = f'{WORKDIR}/{TASK_ID}_initial.pptx'


def create_initial():
    prs = Presentation()
    # Standard widescreen slide dimensions
    # Default is 10 x 7.5 inches

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    layouts = prs.slide_layouts
    # 0=Title Slide, 1=Title+Content, 2=Title+Two Content, 5=Blank, 6=Title Only

    # --- Slide 1: Title Slide ---
    slide = prs.slides.add_slide(layouts[0])
    slide.shapes.title.text = "Computational Neuroscience Lab\nAnnual Results Presentation"
    slide.placeholders[1].text = "Dr. Elena Vasquez — Spring 2025\nInstitute for Brain-Machine Interfaces"

    # --- Slide 2: Agenda ---
    slide = prs.slides.add_slide(layouts[1])
    slide.shapes.title.text = "Agenda"
    tf = slide.placeholders[1].text_frame
    tf.text = "1. Introduction & Background"
    items = [
        "2. Research Objectives",
        "3. Experimental Design",
        "4. Data Collection Methods",
        "5. Results: Electrophysiology",
        "6. Results: Behavioral Metrics",
        "7. Statistical Analysis",
        "8. Key Findings",
        "9. Discussion",
        "10. Future Directions",
        "11. Acknowledgements",
    ]
    for item in items:
        para = tf.add_paragraph()
        para.text = item
        para.level = 0

    # --- Slide 3: Introduction ---
    slide = prs.slides.add_slide(layouts[1])
    slide.shapes.title.text = "Introduction & Background"
    tf = slide.placeholders[1].text_frame
    tf.text = "Neural plasticity in adult cortex following sensory deprivation"
    para = tf.add_paragraph()
    para.text = "Previous work (Zhang et al., 2021) showed cortical remapping over 4-week period"
    para = tf.add_paragraph()
    para.text = "Open questions remain regarding timescale and molecular mechanisms"
    para = tf.add_paragraph()
    para.text = "Our lab developed novel calcium imaging pipeline (NIH R01-NS112345)"

    # --- Slide 4: Research Objectives ---
    slide = prs.slides.add_slide(layouts[1])
    slide.shapes.title.text = "Research Objectives"
    tf = slide.placeholders[1].text_frame
    tf.text = "Primary Aims:"
    objectives = [
        "Aim 1: Characterize temporal dynamics of cortical remapping (n=24 mice)",
        "Aim 2: Identify critical periods for plasticity induction",
        "Aim 3: Map transcriptomic changes in deprived cortex",
        "Aim 4: Test pharmacological interventions to enhance plasticity",
    ]
    for obj in objectives:
        para = tf.add_paragraph()
        para.text = obj
        para.level = 1

    # --- Slide 5: Experimental Design ---
    slide = prs.slides.add_slide(layouts[1])
    slide.shapes.title.text = "Experimental Design"
    tf = slide.placeholders[1].text_frame
    tf.text = "Subjects: C57BL/6J mice, 8–12 weeks, mixed sex (n=48 total)"
    items = [
        "Groups: Sham (n=12), Whisker Deprivation 7d (n=12), WD 14d (n=12), WD 21d (n=12)",
        "Imaging: Two-photon calcium imaging with GCaMP7f (Janelia AAV-hSyn-GCaMP7f)",
        "Region: Barrel cortex (S1BF), layer 2/3, 512×512 px, 30 Hz",
        "Behavioral assay: Texture discrimination task, 5-day training protocol",
    ]
    for item in items:
        para = tf.add_paragraph()
        para.text = item
        para.level = 1

    # --- Slide 6: Data Collection Methods ---
    slide = prs.slides.add_slide(layouts[1])
    slide.shapes.title.text = "Data Collection Methods"
    tf = slide.placeholders[1].text_frame
    tf.text = "Imaging Sessions:"
    items = [
        "Baseline: 3 sessions pre-deprivation (days -7, -3, 0)",
        "Post-deprivation: sessions at days 3, 7, 14, 21",
        "Field of view: 400 × 400 µm, ~200 neurons tracked per session",
        "Motion correction: NoRMCorre algorithm (Pnevmatikakis & Giovannucci, 2017)",
        "Signal extraction: CNMF-E with manual curation (n_pass = 2)",
    ]
    for item in items:
        para = tf.add_paragraph()
        para.text = item
        para.level = 1

    # --- Slide 7: Results - Electrophysiology ---
    slide = prs.slides.add_slide(layouts[1])
    slide.shapes.title.text = "Results: Electrophysiology"
    tf = slide.placeholders[1].text_frame
    tf.text = "Single-unit recordings in barrel cortex:"
    items = [
        "Mean firing rate increased 34% in deprived cortex at 14d (p < 0.001, paired t-test)",
        "Receptive field expansion: principal whisker → adjacent whisker responses",
        "Expansion first significant at day 7 (Δ = 0.42 ± 0.09 spikes/s, p = 0.003)",
        "Strongest effect in supra-granular layers (L2/3) vs. infra-granular (L5/6)",
        "Inter-animal variability: SD = 0.18, CV = 0.22",
    ]
    for item in items:
        para = tf.add_paragraph()
        para.text = item
        para.level = 1

    # --- Slide 8: Results - Calcium Imaging ---
    slide = prs.slides.add_slide(layouts[1])
    slide.shapes.title.text = "Results: Calcium Imaging"
    tf = slide.placeholders[1].text_frame
    tf.text = "Population-level activity changes (two-photon imaging):"
    items = [
        "Total active neurons: 187 ± 23 per FOV at baseline; 204 ± 19 at 14d",
        "Mean ΔF/F increased 18% in deprived barrel column (B2) at 14d",
        "Cross-correlation between B2 and adjacent D2 column increased from 0.31 to 0.54",
        "Principal component analysis: first 3 PCs explain 67% variance at baseline, 71% at 14d",
        "Dimensionality reduction (t-SNE): clear separation of deprivation time points",
    ]
    for item in items:
        para = tf.add_paragraph()
        para.text = item
        para.level = 1

    # --- Slide 9: Results - Behavioral Metrics ---
    slide = prs.slides.add_slide(layouts[1])
    slide.shapes.title.text = "Results: Behavioral Metrics"
    tf = slide.placeholders[1].text_frame
    tf.text = "Texture discrimination task performance:"
    items = [
        "Sham group: 82 ± 4% correct at asymptote (days 3–5 post-training)",
        "WD 7d group: 76 ± 6% correct (not significantly different from sham, p = 0.14)",
        "WD 14d group: 68 ± 8% correct (significant impairment, p = 0.006)",
        "WD 21d group: 65 ± 7% correct (significant impairment, p = 0.002)",
        "Reaction times: no significant group differences (F(3,44) = 1.32, p = 0.28)",
    ]
    for item in items:
        para = tf.add_paragraph()
        para.text = item
        para.level = 1

    # --- Slide 10: Statistical Analysis ---
    slide = prs.slides.add_slide(layouts[1])
    slide.shapes.title.text = "Statistical Analysis"
    tf = slide.placeholders[1].text_frame
    tf.text = "Methods:"
    items = [
        "Primary comparisons: two-way ANOVA (group × time), post-hoc Tukey HSD",
        "Calcium data: linear mixed-effects model (lme4 in R 4.3.1)",
        "Random effects: animal ID, FOV identity",
        "Multiple comparisons: FDR correction (Benjamini-Hochberg, α = 0.05)",
        "Power analysis: 80% power to detect d = 0.6 with n = 12/group (GPower 3.1)",
        "Outlier removal: Grubbs test, 2 animals excluded (hardware failure)",
    ]
    for item in items:
        para = tf.add_paragraph()
        para.text = item
        para.level = 1

    # --- Slide 11: Key Findings ---
    slide = prs.slides.add_slide(layouts[1])
    slide.shapes.title.text = "Key Findings"
    tf = slide.placeholders[1].text_frame
    tf.text = "Summary of main results:"
    items = [
        "Finding 1: Cortical remapping begins within 7 days of whisker deprivation",
        "Finding 2: Behavioral impairment emerges later (14 days) than neural changes",
        "Finding 3: Plasticity magnitude correlates with baseline firing rate (r = 0.61, p = 0.002)",
        "Finding 4: Layer 2/3 shows stronger remapping than Layer 5 (ratio 2.3:1)",
        "Finding 5: MD treatment (10 mg/kg i.p.) partially rescued behavioral deficit at 14d",
    ]
    for item in items:
        para = tf.add_paragraph()
        para.text = item
        para.level = 1

    # --- Slide 12: Discussion ---
    slide = prs.slides.add_slide(layouts[1])
    slide.shapes.title.text = "Discussion"
    tf = slide.placeholders[1].text_frame
    tf.text = "Interpretation of findings:"
    items = [
        "Temporal dissociation supports 'silent remapping' hypothesis (Merzenich, 1983)",
        "Layer-specific effects consistent with thalamocortical vs. corticocortical inputs",
        "Correlation between baseline activity and plasticity suggests homeostatic regulation",
        "Pharmacological rescue (MD) implicates NMDA-R dependent mechanisms",
        "Limitations: single cortical region, one deprivation modality, adult mice only",
    ]
    for item in items:
        para = tf.add_paragraph()
        para.text = item
        para.level = 1

    # --- Slide 13: Comparison with Literature ---
    slide = prs.slides.add_slide(layouts[1])
    slide.shapes.title.text = "Comparison with Published Literature"
    tf = slide.placeholders[1].text_frame
    tf.text = "Our results vs. existing studies:"
    items = [
        "Agreement with Feldman & Brecht (2005): early neural then behavioral effects",
        "Discrepancy with Knott et al. (2002): we found stronger L2/3 than L4 effects",
        "Consistent with Bhatt et al. (2009): spine dynamics track remapping",
        "Novel contribution: MD rescue effect not previously shown in this paradigm",
        "Larger sample size (n=48) than most prior studies (median n=12 in literature)",
    ]
    for item in items:
        para = tf.add_paragraph()
        para.text = item
        para.level = 1

    # --- Slide 14: Future Directions ---
    slide = prs.slides.add_slide(layouts[1])
    slide.shapes.title.text = "Future Directions"
    tf = slide.placeholders[1].text_frame
    tf.text = "Planned experiments (2025–2026):"
    items = [
        "Aim 5: Single-cell RNA-seq of deprived vs. spared barrel columns (10x Genomics)",
        "Aim 6: Optogenetic silencing of L2/3 to test causal role in remapping",
        "Aim 7: Longitudinal imaging through recovery (re-whisker after 21d deprivation)",
        "Collaboration: joint paper with Chen lab (MIT) on cross-modal plasticity",
        "Grant submission: NSF NeuroNex proposal (due Sept 2025)",
    ]
    for item in items:
        para = tf.add_paragraph()
        para.text = item
        para.level = 1

    # --- Slide 15: Acknowledgements ---
    slide = prs.slides.add_slide(layouts[1])
    slide.shapes.title.text = "Acknowledgements"
    tf = slide.placeholders[1].text_frame
    tf.text = "Lab Members:"
    items = [
        "Dr. Marcus Okafor (postdoc) — behavioral experiments",
        "Priya Nair (PhD student) — imaging data collection & analysis",
        "James Whitfield (research technician) — animal husbandry & surgery",
        "Funding: NIH R01-NS112345, NSF CAREER Award 2142876",
        "Equipment: Two-photon microscope supported by S10 OD025132",
        "Computing: XSEDE allocation TG-IBN200015",
    ]
    for item in items:
        para = tf.add_paragraph()
        para.text = item
        para.level = 1

    # --- Slide 16: Questions ---
    slide = prs.slides.add_slide(layouts[0])
    slide.shapes.title.text = "Questions & Discussion"
    slide.placeholders[1].text = "Thank you!\n\nContact: e.vasquez@neurosci.edu\nLab website: vasquezlab.neurosci.edu\nData/code: github.com/vasquezlab/cortical-plasticity-2025"

    # Save to workdir as .pptx (for tracking)
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Also save to the task-specific location in ~/Documents/
    os.makedirs('/home/user/Documents', exist_ok=True)
    import shutil
    shutil.copy(OUTPUT, TASK_FILE)
    print(f'Task file copied to: {TASK_FILE}')
    print(f'Total slides: {len(prs.slides)}')

create_initial()
