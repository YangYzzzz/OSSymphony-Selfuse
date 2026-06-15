"""
Reward Script: Add watermark to all slides in lab_results_talk.odp
Task ID: impress_cross_acad_028
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.4 pts): All 16 slides have a text box containing 'PRELIMINARY - Do Not Distribute'
  - Component 2 (0.3 pts): Watermark has correct font size (24pt) and color (#CCCCCC) on all slides
  - Component 3 (0.3 pts): Watermark is rotated 45 degrees on all slides
"""

import os

from pptx import Presentation
from pptx.util import Pt

WORKDIR = '/home/user'
TASK_ID = 'impress_cross_acad_028'
WATERMARK_TEXT = 'PRELIMINARY - Do Not Distribute'
EXPECTED_SLIDES = 16
EXPECTED_FONT_SIZE_PT = 24
EXPECTED_COLOR = 'CCCCCC'
# 45 degrees = 45 * 60000 = 2700000 in OOXML rotation units
EXPECTED_ROTATION = 2700000
# Allow some tolerance (e.g., +-1 degree = +-60000 units)
ROTATION_TOLERANCE = 60000


def find_watermark_shape(slide):
    """Find the watermark text shape on a slide. Returns shape or None."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if WATERMARK_TEXT in text:
                return shape
    return None


def get_shape_rotation(shape):
    """
    Get rotation of a shape in OOXML units (60000 per degree).
    Returns None if no rotation attribute found.
    """
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    sp_element = shape.element
    for elem in sp_element.iter():
        tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
        if tag == 'xfrm':
            rot_str = elem.attrib.get('rot', None)
            if rot_str is not None:
                return int(rot_str)
            else:
                return 0  # No rotation attribute means 0 degrees
    return 0


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Sanity check: verify slide count
    num_slides = len(prs.slides)
    if num_slides != EXPECTED_SLIDES:
        print(f"WARNING: Expected {EXPECTED_SLIDES} slides, found {num_slides}")

    # Component 1: All 16 slides have a watermark text box (0.4 points)
    # A watermark shape must exist on every slide with the text 'PRELIMINARY - Do Not Distribute'
    try:
        slides_with_watermark = []
        slides_without_watermark = []
        for idx, slide in enumerate(prs.slides):
            wm_shape = find_watermark_shape(slide)
            if wm_shape is not None:
                slides_with_watermark.append(idx + 1)
            else:
                slides_without_watermark.append(idx + 1)

        if len(slides_with_watermark) == EXPECTED_SLIDES:
            print(f"PASS: Component 1 — All {EXPECTED_SLIDES} slides have watermark text box (0.4 pts)")
            total_score += 0.4
        elif len(slides_with_watermark) > 0:
            # Partial: some slides have watermark
            partial = 0.4 * len(slides_with_watermark) / EXPECTED_SLIDES
            print(f"PARTIAL: Component 1 — {len(slides_with_watermark)}/{EXPECTED_SLIDES} slides have watermark "
                  f"(+{partial:.3f} pts). Missing slides: {slides_without_watermark}")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — No slides have watermark text box")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")
        slides_with_watermark = []

    if not slides_with_watermark:
        # No watermarks at all: components 2 and 3 cannot pass either
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {min(total_score, 1.0)}")
        return min(total_score, 1.0)

    # Component 2: Watermark has correct font size (24pt) and color (#CCCCCC) on all slides (0.3 points)
    # Both conditions must be satisfied on every slide that has a watermark
    try:
        slides_correct_style = []
        slides_wrong_style = []

        for idx, slide in enumerate(prs.slides):
            wm_shape = find_watermark_shape(slide)
            if wm_shape is None:
                slides_wrong_style.append(idx + 1)
                continue

            # Check font size and color for runs in watermark
            slide_ok = True
            for para in wm_shape.text_frame.paragraphs:
                for run in para.runs:
                    if not (run.text or '').strip():
                        continue
                    # Check font size
                    font_size_ok = False
                    if run.font.size is not None:
                        # run.font.size is in EMU: 1pt = 12700 EMU
                        size_pt = run.font.size / 12700
                        if abs(size_pt - EXPECTED_FONT_SIZE_PT) < 0.5:
                            font_size_ok = True
                    else:
                        # Try sz attribute in rPr (sz is in hundredths of a point)
                        try:
                            rPr = run.font._element
                            sz = rPr.get('sz')
                            if sz is not None and abs(int(sz) / 100 - EXPECTED_FONT_SIZE_PT) < 0.5:
                                font_size_ok = True
                        except Exception:
                            pass

                    # Check font color
                    color_ok = False
                    try:
                        if run.font.color.type is not None:
                            color_str = str(run.font.color.rgb).upper()
                            if color_str == EXPECTED_COLOR:
                                color_ok = True
                    except Exception:
                        pass

                    if not (font_size_ok and color_ok):
                        slide_ok = False
                        if idx == 0:
                            print(f"  DEBUG Slide 1: font_size_ok={font_size_ok} (size={run.font.size}), "
                                  f"color_ok={color_ok}")
                        break

            if slide_ok:
                slides_correct_style.append(idx + 1)
            else:
                slides_wrong_style.append(idx + 1)

        if len(slides_correct_style) == EXPECTED_SLIDES:
            print(f"PASS: Component 2 — All {EXPECTED_SLIDES} slides have correct font size "
                  f"({EXPECTED_FONT_SIZE_PT}pt) and color (#{EXPECTED_COLOR}) (0.3 pts)")
            total_score += 0.3
        elif len(slides_correct_style) > 0:
            partial = 0.3 * len(slides_correct_style) / EXPECTED_SLIDES
            print(f"PARTIAL: Component 2 — {len(slides_correct_style)}/{EXPECTED_SLIDES} slides have correct "
                  f"font style (+{partial:.3f} pts). Wrong style on slides: {slides_wrong_style[:5]}")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No slides have correct font size/color watermark")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Watermark is rotated 45 degrees on all slides (0.3 points)
    # 45 degrees = 2700000 in OOXML units (60000 per degree)
    try:
        slides_correct_rotation = []
        slides_wrong_rotation = []

        for idx, slide in enumerate(prs.slides):
            wm_shape = find_watermark_shape(slide)
            if wm_shape is None:
                slides_wrong_rotation.append(idx + 1)
                continue

            rotation = get_shape_rotation(wm_shape)
            if abs(rotation - EXPECTED_ROTATION) <= ROTATION_TOLERANCE:
                slides_correct_rotation.append(idx + 1)
            else:
                slides_wrong_rotation.append(idx + 1)
                if idx == 0:
                    print(f"  DEBUG Slide 1: rotation={rotation} (expected ~{EXPECTED_ROTATION}, "
                          f"i.e. {EXPECTED_ROTATION/60000:.1f} degrees)")

        if len(slides_correct_rotation) == EXPECTED_SLIDES:
            print(f"PASS: Component 3 — All {EXPECTED_SLIDES} slides have watermark rotated 45 degrees (0.3 pts)")
            total_score += 0.3
        elif len(slides_correct_rotation) > 0:
            partial = 0.3 * len(slides_correct_rotation) / EXPECTED_SLIDES
            print(f"PARTIAL: Component 3 — {len(slides_correct_rotation)}/{EXPECTED_SLIDES} slides have correct "
                  f"rotation (+{partial:.3f} pts). Wrong rotation on slides: {slides_wrong_rotation[:5]}")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No slides have watermark rotated 45 degrees. "
                  f"Sample rotation: {get_shape_rotation(find_watermark_shape(prs.slides[0])) if find_watermark_shape(prs.slides[0]) else 'N/A'}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against golden file (path on VM)
file_path = f'{WORKDIR}/{TASK_ID}_initial.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
