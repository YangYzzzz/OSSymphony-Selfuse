"""
FINAL REWARD SCRIPT - SUCCESS
Task: Insert one row above row 2 in Table 1.
Generated: 2025-10-17 12:16:21
Status: success
Model: azure-o3
Total Steps: 2
"""

from pptx import Presentation
import os

def verify_table_row_insertion(file_path="/home/user/insert_one_row_above_row_2_in_table_1.pptx"):
    """Verify that exactly one blank row was inserted above the original row 2 in Table 1.

    Scoring (progressive):
        0.4 – Table 1 now has at least one extra row (≥4 rows)
        0.3 – Newly-inserted first row is completely blank
        0.3 – Second row contains the data that used to be the first data row (cell[0] == 'R1C1')
        Total possible: 1.0
    """
    print(f"Checking presentation: {file_path}")

    total_score = 0.0
    max_score = 1.0

    # --------------- Load presentation ---------------
    if not os.path.exists(file_path):
        print("✗ File not found")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Failed to load PPTX: {e}")
        return 0.0

    # --------------- Locate Table 1 ---------------
    target_table = None
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table and shape.name.strip().lower() == "table 1":
                target_table = shape.table
                break
        if target_table is not None:
            break

    if target_table is None:
        print("✗ Could not find a table named 'Table 1'")
        return 0.0

    row_count = len(target_table.rows)
    col_count = len(target_table.columns)
    print(f"Table found with {row_count} rows and {col_count} columns")

    # --------------- Requirement 1: Row count increased ---------------
    # Original table had 3 rows of data. After insertion, should be at least 4.
    if row_count >= 4:
        print("✓ Table has at least 4 rows (row insertion likely performed) (+0.4)")
        total_score += 0.4
    else:
        print("✗ Table does not have enough rows (0 points)")

    # --------------- Requirement 2: First row is blank ---------------
    first_row_blank = True
    for cell in target_table.rows[0].cells:
        txt = cell.text_frame.text if cell.text_frame else ""
        if txt.strip():
            first_row_blank = False
            break
    if first_row_blank:
        print("✓ First row is blank (indicates new row insertion) (+0.3)")
        total_score += 0.3
    else:
        print("✗ First row is not blank (0 points)")

    # --------------- Requirement 3: Original data shifted down ---------------
    expected_value = "r1c1"
    second_row_first_cell = target_table.rows[1].cells[0]
    second_row_text = second_row_first_cell.text_frame.text.strip() if second_row_first_cell.text_frame else ""
    if second_row_text.lower() == expected_value:
        print("✓ Second row contains original first row data (R1C1) (+0.3)")
        total_score += 0.3
    else:
        print(f"✗ Second row first cell unexpected value: '{second_row_text}' (0 points)")

    # --------------- Final score ---------------
    final_score = min(total_score, max_score)
    print(f"Total Score: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score

# Execute verification when script is run
if __name__ == "__main__":
    verify_table_row_insertion()

