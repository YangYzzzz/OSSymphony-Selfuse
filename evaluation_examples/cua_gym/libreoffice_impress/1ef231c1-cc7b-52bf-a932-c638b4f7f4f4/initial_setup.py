"""
Initial Setup: Create a 15-slide product launch presentation
Task ID: impress_gf5_013
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_013'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            body.paragraphs[0].text = bullet
        else:
            p = body.add_paragraph()
            p.text = bullet
            p.level = 0
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide (Executive + Technical)
    add_title_slide(prs, "NovaTech SmartHub Pro Launch",
                    "Product Launch Strategy & Technical Architecture\nQ3 2025")

    # Slide 2: Executive Summary (Executive)
    add_content_slide(prs, "Executive Summary", [
        "SmartHub Pro represents a $4.2M market opportunity in IoT home automation",
        "Target launch date: September 15, 2025",
        "Projected first-year revenue: $12.8M with 23% margin",
        "Key differentiator: AI-powered predictive home management",
        "Strategic partnerships with 3 major retailers confirmed",
    ])

    # Slide 3: Technical Architecture Overview (Technical)
    add_content_slide(prs, "Technical Architecture Overview", [
        "Microservices-based backend on AWS EKS (Kubernetes)",
        "Real-time data pipeline: Kafka → Flink → TimescaleDB",
        "Edge computing module: ARM Cortex-M7 with 512KB SRAM",
        "Communication protocols: Zigbee 3.0, Thread, Wi-Fi 6E",
        "End-to-end encryption with AES-256-GCM",
    ])

    # Slide 4: Hardware Specifications (Technical)
    add_content_slide(prs, "Hardware Specifications", [
        "Processor: Quad-core ARM Cortex-A55 @ 2.0 GHz",
        "Memory: 4GB LPDDR4X RAM, 32GB eMMC storage",
        "Sensors: Temperature, humidity, ambient light, motion (PIR)",
        "Connectivity: Wi-Fi 6E, Bluetooth 5.3, Zigbee 3.0, Thread",
        "Power: 5V/3A USB-C, optional PoE+ (IEEE 802.3at)",
        "Dimensions: 142mm × 142mm × 38mm, weight 285g",
    ])

    # Slide 5: Market Analysis (Executive)
    add_content_slide(prs, "Market Analysis & Opportunity", [
        "Global smart home market projected at $338B by 2027 (CAGR 14.2%)",
        "Target segment: Premium homeowners, age 30-55, household income >$120K",
        "Competitor gap: No unified hub with predictive AI at sub-$300 price point",
        "Early adopter survey: 78% expressed strong purchase intent",
        "Regional focus: North America (60%), Europe (25%), APAC (15%)",
    ])

    # Slide 6: API & Integration Framework (Technical)
    add_content_slide(prs, "API & Integration Framework", [
        "RESTful API v2.0 with OpenAPI 3.1 specification",
        "GraphQL endpoint for real-time device state queries",
        "SDK available for Python, JavaScript, Swift, and Kotlin",
        "IFTTT, HomeAssistant, and SmartThings integrations",
        "Webhook support for custom automation workflows",
        "Rate limiting: 1000 req/min per device, 10000 req/min per account",
    ])

    # Slide 7: Security & Compliance (Technical)
    add_content_slide(prs, "Security Architecture & Compliance", [
        "Zero-trust architecture with mutual TLS between all services",
        "FIDO2/WebAuthn for passwordless device authentication",
        "SOC 2 Type II certification in progress (ETA: August 2025)",
        "GDPR, CCPA, and PIPEDA compliant data handling",
        "Automated vulnerability scanning via Snyk and Dependabot",
        "Bug bounty program launching with HackerOne",
    ])

    # Slide 8: Financial Projections (Executive)
    add_content_slide(prs, "Financial Projections (FY2025-2027)", [
        "FY2025 Revenue: $12.8M | COGS: $7.4M | Gross Margin: 42%",
        "FY2026 Revenue: $28.5M | COGS: $15.2M | Gross Margin: 47%",
        "FY2027 Revenue: $51.3M | COGS: $25.1M | Gross Margin: 51%",
        "Break-even expected by Q2 2026 at 45,000 units sold",
        "R&D investment: $3.2M in Year 1, scaling to $5.8M by Year 3",
        "Customer acquisition cost: $42 per unit, LTV: $380",
    ])

    # Slide 9: Testing & QA Pipeline (Technical)
    add_content_slide(prs, "Testing & Quality Assurance Pipeline", [
        "Automated CI/CD via GitHub Actions with 94% code coverage",
        "Hardware-in-the-loop (HIL) testing with 200-device test farm",
        "Performance benchmarks: <50ms response time for local commands",
        "Stress testing: 10,000 concurrent device connections sustained",
        "UL 2900-1 cybersecurity certification testing underway",
        "Beta testing program: 500 households across 12 cities",
    ])

    # Slide 10: Deployment & Scalability (Technical)
    add_content_slide(prs, "Deployment & Scalability Plan", [
        "Blue-green deployment strategy with automated rollback",
        "Multi-region: us-east-1, eu-west-1, ap-southeast-1",
        "Auto-scaling: 10K to 1M connected devices without config changes",
        "OTA firmware updates with A/B partition scheme",
        "CDN for static assets via CloudFront (99.99% SLA)",
        "Database sharding strategy for >10M device telemetry records/day",
    ])

    # Slide 11: Manufacturing & Supply Chain
    add_content_slide(prs, "Manufacturing & Supply Chain", [
        "Primary manufacturer: Foxconn (Shenzhen facility)",
        "Secondary source: Flex Ltd (Guadalajara) for risk mitigation",
        "Lead time: 8-10 weeks from PO to warehouse delivery",
        "Initial production run: 25,000 units (August 2025)",
        "Component inventory buffer: 6 weeks of safety stock",
        "Landed cost per unit: $164 (target: $148 at scale by Q2 2026)",
    ])

    # Slide 12: Go-to-Market Strategy (Executive)
    add_content_slide(prs, "Go-to-Market Strategy", [
        "Launch event: September 15, 2025 at CES Unveiled New York",
        "Retail partners: Best Buy, Amazon, Home Depot (confirmed)",
        "Digital campaign: $1.8M budget across Google, Meta, and YouTube",
        "Influencer program: 50 tech reviewers with early access units",
        "Pricing: $279 MSRP, $249 early-bird (first 5,000 units)",
        "30-day money-back guarantee + 2-year limited warranty",
    ])

    # Slide 13: Customer Support Plan
    add_content_slide(prs, "Customer Support & Service Plan", [
        "24/7 in-app chat support with AI-assisted triage",
        "Phone support: 8am-10pm EST, 7 days/week",
        "Target first-response time: <2 minutes for chat, <5 minutes for phone",
        "Self-service knowledge base with 200+ articles at launch",
        "Community forum with gamified expert user program",
        "NPS target: >60 within first 6 months post-launch",
    ])

    # Slide 14: Risk Assessment
    add_content_slide(prs, "Risk Assessment & Mitigation", [
        "Supply chain disruption: Dual-source strategy + 6-week buffer",
        "Competitor response: 6-month feature lead with patent portfolio",
        "Regulatory changes: Legal team monitoring IoT legislation in 12 jurisdictions",
        "Cybersecurity breach: $5M cyber insurance + incident response plan",
        "Adoption slower than projected: Flexible marketing budget reallocation",
    ])

    # Slide 15: Timeline & Milestones
    add_content_slide(prs, "Launch Timeline & Key Milestones", [
        "June 2025: Final hardware design freeze",
        "July 2025: FCC/CE certification submission",
        "August 2025: Initial production run (25,000 units)",
        "September 1, 2025: Retail partner inventory delivery",
        "September 15, 2025: Public launch event",
        "October 2025: Post-launch firmware update (v1.1)",
        "December 2025: International expansion (UK, Germany, Japan)",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
