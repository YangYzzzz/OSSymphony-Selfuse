"""
Initial Setup: Create undergraduate presentation with 12 slides (no review comments)
Task ID: impress_cross_acad_034
Domain: libreoffice_impress
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_cross_acad_034'
OUTPUT = f'{WORKDIR}/Documents/{TASK_ID}_initial.pptx'
# Also place at the path referenced in the task instruction
TASK_PATH = f'{WORKDIR}/Documents/undergrad_presentation.odp'

def add_title_slide(prs, title, subtitle):
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title, content_lines):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.text = content_lines[0]
    for line in content_lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 1
    return slide

def add_section_header(prs, title, subtitle=''):
    layout = prs.slide_layouts[2] if len(prs.slide_layouts) > 2 else prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        try:
            slide.placeholders[1].text = subtitle
        except Exception:
            pass
    return slide

def add_text_box(slide, left_cm, top_cm, width_cm, height_cm, text, font_size=14, bold=False, color=None):
    txBox = slide.shapes.add_textbox(
        Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)
    return txBox

def create_initial():
    prs = Presentation()
    # Standard widescreen (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    layout0 = prs.slide_layouts[0]
    s1 = prs.slides.add_slide(layout0)
    s1.shapes.title.text = "Impact of Social Media on Academic Performance"
    try:
        s1.placeholders[1].text = "A Quantitative Study\nEmily Hartwell\nDepartment of Psychology\nUniversity of Westfield, 2025"
    except Exception:
        pass

    # Slide 2: Introduction / Research Background
    layout1 = prs.slide_layouts[1]
    s2 = prs.slides.add_slide(layout1)
    s2.shapes.title.text = "Introduction"
    tf2 = s2.placeholders[1].text_frame
    tf2.text = "Background"
    for line in [
        "Social media platforms have become ubiquitous in college students' lives",
        "Average student spends 4.2 hours/day on social media (Pew Research, 2024)",
        "Conflicting literature on academic impact",
        "Gap: limited longitudinal studies on GPA correlation",
    ]:
        p = tf2.add_paragraph()
        p.text = line
        p.level = 1

    # Slide 3: Literature Review (will get comment "Needs a citation here")
    s3 = prs.slides.add_slide(layout1)
    s3.shapes.title.text = "Literature Review"
    tf3 = s3.placeholders[1].text_frame
    tf3.text = "Key Prior Studies"
    for line in [
        "Positive association: social media for academic collaboration increases engagement",
        "Negative association: passive scrolling correlated with lower GPA (r = -0.31)",
        "Mixed findings: platform type matters — LinkedIn vs. TikTok differ significantly",
        "Most studies limited to single-semester cross-sectional designs",
        "Social media usage peaks during exam periods among undergraduates",
    ]:
        p = tf3.add_paragraph()
        p.text = line
        p.level = 1

    # Slide 4: Research Questions
    s4 = prs.slides.add_slide(layout1)
    s4.shapes.title.text = "Research Questions"
    tf4 = s4.placeholders[1].text_frame
    tf4.text = "This study addresses the following questions:"
    questions = [
        "RQ1: Does daily social media usage time negatively predict semester GPA?",
        "RQ2: Does the type of platform moderate the relationship?",
        "RQ3: Are these effects consistent across academic disciplines?",
    ]
    for q in questions:
        p = tf4.add_paragraph()
        p.text = q
        p.level = 1

    # Slide 5: Methodology - Participants (will get comment "This figure is too small")
    s5 = prs.slides.add_slide(layout1)
    s5.shapes.title.text = "Methodology: Participants & Data"
    tf5 = s5.placeholders[1].text_frame
    tf5.text = "Sample"
    for line in [
        "N = 412 undergraduate students across 4 universities",
        "Age range: 18–24 (M = 20.3, SD = 1.7)",
        "Gender: 54% female, 44% male, 2% non-binary",
        "Disciplines: STEM (38%), Humanities (29%), Social Sciences (33%)",
    ]:
        p = tf5.add_paragraph()
        p.text = line
        p.level = 1
    # Add a placeholder figure description (the "figure" that's "too small")
    fig_box = s5.shapes.add_textbox(Cm(18), Cm(3), Cm(5), Cm(4))
    fig_tf = fig_box.text_frame
    fig_tf.word_wrap = True
    fig_p = fig_tf.paragraphs[0]
    fig_p.text = "[Fig. 1: Sample Distribution by Discipline]"
    fig_p.alignment = PP_ALIGN.CENTER
    fig_run = fig_p.runs[0]
    fig_run.font.size = Pt(9)
    fig_run.font.italic = True
    fig_run.font.color.rgb = RGBColor(0x60, 0x60, 0x60)

    # Slide 6: Measures / Instruments
    s6 = prs.slides.add_slide(layout1)
    s6.shapes.title.text = "Measures & Instruments"
    tf6 = s6.placeholders[1].text_frame
    tf6.text = "Data Collection Tools"
    for line in [
        "Social Media Usage Scale (SMUS-12): validated 12-item Likert instrument",
        "Academic GPA: obtained from registrar records with consent",
        "Platform Usage Log: self-reported daily minutes per platform",
        "Academic Engagement Survey (AES): 8-item scale (α = 0.87)",
        "Control variables: part-time employment, extracurriculars, sleep hours",
    ]:
        p = tf6.add_paragraph()
        p.text = line
        p.level = 1

    # Slide 7: Results - Descriptive Statistics
    s7 = prs.slides.add_slide(layout1)
    s7.shapes.title.text = "Results: Descriptive Statistics"
    tf7 = s7.placeholders[1].text_frame
    tf7.text = "Key Descriptive Findings"
    for line in [
        "Mean daily social media use: 3.8 hrs (SD = 1.6)",
        "Mean cumulative GPA: 3.14 (SD = 0.52)",
        "Most used platforms: Instagram (78%), TikTok (71%), Snapchat (64%)",
        "Academic collaboration use: 42% use platforms for study groups",
        "Peak usage: 11pm–1am reported by 67% of participants",
    ]:
        p = tf7.add_paragraph()
        p.text = line
        p.level = 1

    # Slide 8: Results - Regression Analysis (will get comment "Excellent visualization!")
    s8 = prs.slides.add_slide(layout1)
    s8.shapes.title.text = "Results: Regression Analysis"
    tf8 = s8.placeholders[1].text_frame
    tf8.text = "Multiple Regression Model (N=412, R² = 0.41)"
    for line in [
        "Daily usage time: β = -0.28, p < 0.001 (significant negative predictor)",
        "Platform type (TikTok vs. LinkedIn): β = -0.19, p < 0.01",
        "Academic use moderates negative effect: β = +0.23, p < 0.001",
        "Sleep hours: β = +0.31, p < 0.001 (strongest predictor)",
        "Model explains 41% of variance in GPA",
    ]:
        p = tf8.add_paragraph()
        p.text = line
        p.level = 1
    # A descriptive note about the visualization
    viz_box = s8.shapes.add_textbox(Cm(2), Cm(5.5), Cm(20), Cm(1))
    viz_tf = viz_box.text_frame
    viz_p = viz_tf.paragraphs[0]
    viz_p.text = "Figure 2: Path diagram showing standardized regression coefficients"
    viz_run = viz_p.runs[0]
    viz_run.font.size = Pt(10)
    viz_run.font.italic = True
    viz_run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # Slide 9: Discussion
    s9 = prs.slides.add_slide(layout1)
    s9.shapes.title.text = "Discussion"
    tf9 = s9.placeholders[1].text_frame
    tf9.text = "Interpretation of Findings"
    for line in [
        "Passive consumption (scrolling) most harmful to GPA",
        "Active academic use (study groups) buffers negative effects",
        "TikTok associated with larger GPA declines than LinkedIn",
        "Sleep quality mediates relationship between usage and performance",
        "Findings align with displacement hypothesis (Kraut et al., 2002)",
    ]:
        p = tf9.add_paragraph()
        p.text = line
        p.level = 1

    # Slide 10: Limitations
    s10 = prs.slides.add_slide(layout1)
    s10.shapes.title.text = "Limitations"
    tf10 = s10.placeholders[1].text_frame
    tf10.text = "Study Constraints"
    for line in [
        "Self-reported usage data subject to recall bias",
        "Cross-sectional design limits causal inferences",
        "Sample skewed toward large research universities",
        "Unmeasured confounds: family support, financial stress",
        "Rapidly changing platform landscape (algorithm updates)",
    ]:
        p = tf10.add_paragraph()
        p.text = line
        p.level = 1

    # Slide 11: Future Directions
    s11 = prs.slides.add_slide(layout1)
    s11.shapes.title.text = "Future Directions & Implications"
    tf11 = s11.placeholders[1].text_frame
    tf11.text = "Recommendations"
    for line in [
        "Longitudinal design tracking students across all 4 years",
        "Screen-time API integration for objective usage measurement",
        "Intervention study: digital wellness curriculum impact",
        "Practical: universities should promote intentional platform use",
        "Policy: academic social media guidelines for residential programs",
    ]:
        p = tf11.add_paragraph()
        p.text = line
        p.level = 1

    # Slide 12: References / Q&A
    s12 = prs.slides.add_slide(layout1)
    s12.shapes.title.text = "References & Questions"
    tf12 = s12.placeholders[1].text_frame
    tf12.text = "Selected References"
    for line in [
        "Andreassen, C. S. et al. (2016). Development of a Facebook Addiction Scale. Psychological Reports.",
        "Junco, R. (2012). Too much face and not enough books. Computers in Human Behavior.",
        "Pew Research Center (2024). Social Media Use Among College Students.",
        "Twenge, J. M., & Campbell, W. K. (2019). Media use is linked to lower well-being. PsychologicalScience.",
        "",
        "Thank you — Questions?",
    ]:
        p = tf12.add_paragraph()
        p.text = line
        p.level = 0 if not line else 1

    # Ensure Documents directory exists on VM and save
    os.makedirs(f'{WORKDIR}/Documents', exist_ok=True)
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Also copy to the .odp path the task instruction references
    import shutil
    shutil.copy(OUTPUT, TASK_PATH)
    print(f'Also copied to task path: {TASK_PATH}')

create_initial()
