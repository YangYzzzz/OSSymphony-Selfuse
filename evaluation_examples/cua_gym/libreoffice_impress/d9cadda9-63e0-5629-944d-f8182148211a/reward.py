"""
Reward Script: Add review comment text boxes to specific slides in undergrad_presentation.odp
Task ID: impress_cross_acad_034
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 3 has comment text box 'Needs a citation here' (0.4 pts)
    - Sub 1a: text box with correct text exists (0.2 pts)
    - Sub 1b: text box has yellow background #FFFF99 (0.1 pts)
    - Sub 1c: text box positioned at bottom-left ~1cm, ~17cm (0.1 pts)
  Component 2: Slide 5 has comment text box 'This figure is too small' (0.3 pts)
    - Sub 2a: text box with correct text exists (0.15 pts)
    - Sub 2b: text box has yellow background #FFFF99 (0.075 pts)
    - Sub 2c: text box positioned at bottom-left ~1cm, ~17cm (0.075 pts)
  Component 3: Slide 8 has comment text box 'Excellent visualization!' (0.3 pts)
    - Sub 3a: text box with correct text exists (0.15 pts)
    - Sub 3b: text box has yellow background #FFFF99 (0.075 pts)
    - Sub 3c: text box positioned at bottom-left ~1cm, ~17cm (0.075 pts)
  Total: 1.0

Ground truth from task context:
  - Slides 3, 5, 8 should each have an additional text box at bottom-left
  - Position: ~1cm from left, ~17cm from top
  - Yellow background: #FFFF99
  - Black text, 12pt font
  - Comments: 'Needs a citation here', 'This figure is too small', 'Excellent visualization!'
"""

import os
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user/Documents'
TASK_ID = 'impress_cross_acad_034'

YELLOW_BG = 'FFFF99'

# Position thresholds in EMU: ~1cm left = 360000, ~17cm top = 6120000
# Allow 10% position tolerance for bottom-left placement
EXPECTED_LEFT_EMU = 360000
EXPECTED_TOP_EMU = 6120000
POSITION_TOLERANCE = 0.10


def emu_approx(val, expected, tol=POSITION_TOLERANCE):
    """Check if val is within tolerance of expected (relative)."""
    if expected == 0:
        return val == 0
    return abs(val - expected) / max(abs(val), abs(expected)) <= tol


def find_comment_textbox(slide, expected_text):
    """
    Search slide for a TEXT_BOX shape with text matching expected_text.
    Returns dict with detection results for each sub-check.
    """
    result = {
        'text_found': False,
        'yellow_bg': False,
        'position_ok': False,
        'shape_left': None,
        'shape_top': None,
    }

    for shape in slide.shapes:
        # Only check text box shapes (shape_type == 17)
        if shape.shape_type != 17:
            continue
        if not shape.has_text_frame:
            continue

        actual_text = shape.text_frame.text.strip()
        if actual_text != expected_text:
            continue

        # Correct text found
        result['text_found'] = True
        result['shape_left'] = shape.left
        result['shape_top'] = shape.top

        # Check yellow background fill
        try:
            if str(shape.fill.type) == 'SOLID (1)':
                fg_rgb = str(shape.fill.fore_color.rgb)
                if fg_rgb.upper() == YELLOW_BG.upper():
                    result['yellow_bg'] = True
        except Exception:
            pass

        # Check bottom-left position (~1cm left, ~17cm top)
        if emu_approx(shape.left, EXPECTED_LEFT_EMU) and emu_approx(shape.top, EXPECTED_TOP_EMU):
            result['position_ok'] = True

        break  # Stop after finding the first matching text box

    return result


def check_comment_on_slide(prs, slide_index, expected_text, base_pts, slide_num):
    """
    Check if slide at slide_index (0-based) has a comment text box.
    Returns score earned for this component.
    """
    slide = prs.slides[slide_index]
    d = find_comment_textbox(slide, expected_text)

    component_score = 0.0

    # Sub-check A: text box with correct text exists
    if d['text_found']:
        component_score += base_pts * 0.5
        print(f"PASS: Slide {slide_num} — text box with text '{expected_text}' found (+{base_pts * 0.5:.3f} pts)")
    else:
        print(f"FAIL: Slide {slide_num} — no text box with text '{expected_text}' found")
        return 0.0  # If text not found, no further points possible

    # Sub-check B: yellow background
    if d['yellow_bg']:
        component_score += base_pts * 0.25
        print(f"PASS: Slide {slide_num} — text box has yellow background #FFFF99 (+{base_pts * 0.25:.3f} pts)")
    else:
        print(f"FAIL: Slide {slide_num} — text box missing yellow background #FFFF99")

    # Sub-check C: bottom-left position
    if d['position_ok']:
        component_score += base_pts * 0.25
        print(f"PASS: Slide {slide_num} — text box at bottom-left (~1cm left, ~17cm top) (+{base_pts * 0.25:.3f} pts)")
    else:
        left_cm = (d['shape_left'] or 0) / 360000
        top_cm = (d['shape_top'] or 0) / 360000
        print(f"FAIL: Slide {slide_num} — text box position ({left_cm:.1f}cm left, {top_cm:.1f}cm top) not at expected bottom-left (~1cm, ~17cm)")

    return component_score


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition gate: verify slide count is 12
    if len(prs.slides) != 12:
        print(f"CRITICAL: Expected 12 slides, found {len(prs.slides)} — file structure invalid")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide 3 has comment 'Needs a citation here' (0.4 points)
    try:
        c1_score = check_comment_on_slide(prs, slide_index=2, expected_text='Needs a citation here',
                                          base_pts=0.4, slide_num=3)
        if c1_score > 0:
            total_score += c1_score
        print(f"  Component 1 total: {c1_score:.3f}/0.4")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 5 has comment 'This figure is too small' (0.3 points)
    try:
        c2_score = check_comment_on_slide(prs, slide_index=4, expected_text='This figure is too small',
                                          base_pts=0.3, slide_num=5)
        if c2_score > 0:
            total_score += c2_score
        print(f"  Component 2 total: {c2_score:.3f}/0.3")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 8 has comment 'Excellent visualization!' (0.3 points)
    try:
        c3_score = check_comment_on_slide(prs, slide_index=7, expected_text='Excellent visualization!',
                                          base_pts=0.3, slide_num=8)
        if c3_score > 0:
            total_score += c3_score
        print(f"  Component 3 total: {c3_score:.3f}/0.3")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 3)
    print(f"\nScore: {total_score:.3f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against golden file (path on VM)
file_path = f'{WORKDIR}/{TASK_ID}_initial.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
