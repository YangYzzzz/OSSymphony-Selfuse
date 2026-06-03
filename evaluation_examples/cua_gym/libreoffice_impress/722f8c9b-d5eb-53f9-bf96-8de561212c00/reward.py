"""
Reward Script: Change title font to Verdana on slides 1-3
Task ID: osworld_impress_global_font_change_004
Domain: libreoffice_impress

Scoring Rubric:
  Component 1: Slide 1 title font changed to Verdana AND content stays Calibri (0.35 pts)
  Component 2: Slide 2 title font changed to Verdana AND content stays Calibri (0.35 pts)
  Component 3: Slide 3 title font changed to Verdana AND slides 4-6 titles remain
               Times New Roman (unchanged) AND content stays Calibri             (0.30 pts)
  Total: 1.0

Each component FAILS on initial_env (titles were Times New Roman before the task)
and PASSES on golden_env (titles changed to Verdana).
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_global_font_change_004'

# Placeholder type integer values (PP_PLACEHOLDER)
# TITLE = 1, CENTER_TITLE = 3
TITLE_PH_TYPE_VALS = {1, 3}


def get_title_shape(slide):
    """Return the title placeholder shape on the given slide, or None."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            ph = shape.placeholder_format
            if ph is not None and ph.type is not None:
                if int(ph.type) in TITLE_PH_TYPE_VALS:
                    return shape
    return None


def get_all_run_fonts(shape):
    """Return list of (run_text, font_name) for all non-empty runs in a shape."""
    result = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if (run.text or '').strip():
                result.append((run.text, run.font.name))
    return result


def get_content_fonts(slide):
    """Return set of font names from all non-title text shapes on the slide."""
    fonts = set()
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        ph = shape.placeholder_format
        if ph is not None and ph.type is not None:
            if int(ph.type) in TITLE_PH_TYPE_VALS:
                continue  # skip title shapes
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if (run.text or '').strip():
                    fonts.add(run.font.name)
    return fonts


def check_title_is_verdana(slide, slide_num):
    """Return (passed: bool, detail_msg: str)"""
    title_shape = get_title_shape(slide)
    if title_shape is None:
        return False, f"Slide {slide_num}: no title placeholder found"
    runs = get_all_run_fonts(title_shape)
    if not runs:
        return False, f"Slide {slide_num}: title has no non-empty runs"
    non_verdana = [(t, f) for t, f in runs if f != 'Verdana']
    if non_verdana:
        return False, f"Slide {slide_num}: title runs not Verdana: {non_verdana[:3]}"
    return True, f"Slide {slide_num}: all title runs use Verdana"


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Sanity check: must have 6 slides
    if len(prs.slides) != 6:
        print(f"CRITICAL: Expected 6 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)

    # Component 1: Slide 1 title font is Verdana AND slide 1 content stays Calibri (0.35 pts)
    # This FAILS on initial (slide 1 title is Times New Roman) and PASSES on golden (Verdana)
    try:
        passed, msg = check_title_is_verdana(slides[0], slide_num=1)
        if not passed:
            print(f"FAIL: Component 1 — {msg}")
        else:
            # Sub-condition: content on slide 1 stays Calibri
            content_fonts = get_content_fonts(slides[0])
            non_calibri = {f for f in content_fonts if f != 'Calibri'}
            if non_calibri:
                print(f"FAIL: Component 1 — {msg} but slide 1 content changed to: {non_calibri}")
            else:
                print(f"PASS: Component 1 — {msg} and content is Calibri (0.35 pts)")
                total_score += 0.35
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 2 title font is Verdana AND slide 2 content stays Calibri (0.35 pts)
    # This FAILS on initial (slide 2 title is Times New Roman) and PASSES on golden (Verdana)
    try:
        passed, msg = check_title_is_verdana(slides[1], slide_num=2)
        if not passed:
            print(f"FAIL: Component 2 — {msg}")
        else:
            # Sub-condition: content on slide 2 stays Calibri
            content_fonts = get_content_fonts(slides[1])
            non_calibri = {f for f in content_fonts if f != 'Calibri'}
            if non_calibri:
                print(f"FAIL: Component 2 — {msg} but slide 2 content changed to: {non_calibri}")
            else:
                print(f"PASS: Component 2 — {msg} and content is Calibri (0.35 pts)")
                total_score += 0.35
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 3 title font is Verdana AND slides 4-6 title fonts remain Times New Roman
    # AND slide 3 content stays Calibri (0.30 pts)
    # This FAILS on initial (slide 3 title is Times New Roman) and PASSES on golden (Verdana).
    # The compound check anchors the "slides 4-6 unchanged" requirement to the core task change.
    try:
        passed, msg = check_title_is_verdana(slides[2], slide_num=3)
        if not passed:
            print(f"FAIL: Component 3 — {msg}")
        else:
            # Sub-condition: slides 4-6 titles must remain Times New Roman
            unchanged_ok = True
            for slide_idx in [3, 4, 5]:
                title_shape = get_title_shape(slides[slide_idx])
                if title_shape is None:
                    print(f"FAIL: Component 3 — Slide {slide_idx+1} has no title placeholder")
                    unchanged_ok = False
                    break
                runs = get_all_run_fonts(title_shape)
                non_tnr = [(t, f) for t, f in runs if f != 'Times New Roman']
                if non_tnr:
                    print(f"FAIL: Component 3 — Slide {slide_idx+1} title should stay Times New Roman, found: {non_tnr[:2]}")
                    unchanged_ok = False
                    break
            if unchanged_ok:
                # Sub-condition: content on slide 3 stays Calibri
                content_fonts = get_content_fonts(slides[2])
                non_calibri = {f for f in content_fonts if f != 'Calibri'}
                if non_calibri:
                    print(f"FAIL: Component 3 — {msg} but slide 3 content changed to: {non_calibri}")
                else:
                    print(f"PASS: Component 3 — {msg}, slides 4-6 titles remain Times New Roman, slide 3 content is Calibri (0.30 pts)")
                    total_score += 0.30
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Run against canonical artifact path on the VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
