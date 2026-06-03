"""
Initial Setup: 6-slide team onboarding deck with Times New Roman title fonts on all slides.
Task ID: osworld_impress_global_font_change_004
Domain: libreoffice_impress

Creates a 6-slide presentation where ALL title placeholders use Times New Roman.
The task requires the agent to change slides 1-3 title font to Verdana.
Slides 4-6 titles should remain Times New Roman after task completion.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_global_font_change_004'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def apply_title_font(title_shape, font_name, font_size_pt, bold=False):
    """Apply font properties to all runs in the title placeholder."""
    for para in title_shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size_pt)
            run.font.bold = bold


def apply_content_font(content_shape, font_name, font_size_pt):
    """Apply font properties to all runs in a content placeholder."""
    for para in content_shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size_pt)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Layout indices in default python-pptx template:
    #   0 = Title Slide (title + subtitle)
    #   1 = Title and Content
    layout_title = prs.slide_layouts[0]
    layout_content = prs.slide_layouts[1]

    # -----------------------------------------------------------------------
    # Slide 1: Title Slide — Welcome (Times New Roman title — to be changed)
    # -----------------------------------------------------------------------
    slide1 = prs.slides.add_slide(layout_title)
    title1 = slide1.shapes.title
    title1.text = "Welcome to Novara Technologies"
    sub1 = slide1.placeholders[1]
    sub1.text = "Team Onboarding Program — Q2 2025"

    apply_title_font(title1, "Times New Roman", 40, bold=True)
    apply_content_font(sub1, "Calibri", 24)

    # -----------------------------------------------------------------------
    # Slide 2: Company Overview (Times New Roman title — to be changed)
    # -----------------------------------------------------------------------
    slide2 = prs.slides.add_slide(layout_content)
    title2 = slide2.shapes.title
    title2.text = "About Novara Technologies"
    content2 = slide2.placeholders[1]
    content2.text = (
        "Founded in 2010, Novara Technologies is a global leader in enterprise software.\n"
        "Headquarters: Singapore | Offices in 12 countries\n"
        "Employees: 3,400+ worldwide\n"
        "Revenue FY2024: $820M | Growth: 18% YoY"
    )

    apply_title_font(title2, "Times New Roman", 36, bold=False)
    apply_content_font(content2, "Calibri", 18)

    # -----------------------------------------------------------------------
    # Slide 3: Team Structure (Times New Roman title — to be changed)
    # -----------------------------------------------------------------------
    slide3 = prs.slides.add_slide(layout_content)
    title3 = slide3.shapes.title
    title3.text = "Your Team & Reporting Structure"
    content3 = slide3.placeholders[1]
    content3.text = (
        "Engineering Division\n"
        "VP Engineering: Dr. Priya Nair\n"
        "Team Lead: Marcus Webb\n"
        "Senior Engineers: Emily Tan, Jordan Brooks, Sam Patel\n"
        "Junior Engineers: Alex Kim, Lucia Morales\n"
        "Weekly stand-ups: Monday & Thursday, 9:30 AM SGT"
    )

    apply_title_font(title3, "Times New Roman", 36, bold=False)
    apply_content_font(content3, "Calibri", 18)

    # -----------------------------------------------------------------------
    # Slide 4: Benefits & Perks (Times New Roman title — NOT changed by task)
    # -----------------------------------------------------------------------
    slide4 = prs.slides.add_slide(layout_content)
    title4 = slide4.shapes.title
    title4.text = "Benefits & Perks"
    content4 = slide4.placeholders[1]
    content4.text = (
        "Health & Wellness\n"
        "Comprehensive medical, dental, and vision coverage\n"
        "SGD 1,200 annual wellness allowance\n"
        "Mental health support via EAP program\n"
        "Financial Benefits\n"
        "Competitive base salary and performance bonus\n"
        "Stock options (ESOP) — vesting over 4 years\n"
        "Annual salary review in December"
    )

    apply_title_font(title4, "Times New Roman", 36, bold=False)
    apply_content_font(content4, "Calibri", 18)

    # -----------------------------------------------------------------------
    # Slide 5: Policies & Procedures (Times New Roman title — NOT changed)
    # -----------------------------------------------------------------------
    slide5 = prs.slides.add_slide(layout_content)
    title5 = slide5.shapes.title
    title5.text = "Key Policies & Procedures"
    content5 = slide5.placeholders[1]
    content5.text = (
        "Working Hours & Flexibility\n"
        "Core hours: 10 AM to 3 PM SGT\n"
        "Hybrid policy: 3 days in-office minimum\n"
        "20 days annual leave plus 11 public holidays\n"
        "Code of Conduct\n"
        "Zero-tolerance harassment policy\n"
        "Data privacy and confidentiality agreement\n"
        "IT security guidelines — access intranet portal"
    )

    apply_title_font(title5, "Times New Roman", 36, bold=False)
    apply_content_font(content5, "Calibri", 18)

    # -----------------------------------------------------------------------
    # Slide 6: Next Steps (Times New Roman title — NOT changed by task)
    # -----------------------------------------------------------------------
    slide6 = prs.slides.add_slide(layout_content)
    title6 = slide6.shapes.title
    title6.text = "Your First 30 Days"
    content6 = slide6.placeholders[1]
    content6.text = (
        "Week 1: Orientation\n"
        "IT setup, badge access, system accounts\n"
        "Meet your buddy — assigned on Day 1\n"
        "Week 2-3: Role Ramp-Up\n"
        "Shadow senior team members\n"
        "Complete mandatory training modules (LMS)\n"
        "Week 4: Integration\n"
        "Join first sprint planning session\n"
        "30-day check-in with manager"
    )

    apply_title_font(title6, "Times New Roman", 36, bold=False)
    apply_content_font(content6, "Calibri", 18)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
