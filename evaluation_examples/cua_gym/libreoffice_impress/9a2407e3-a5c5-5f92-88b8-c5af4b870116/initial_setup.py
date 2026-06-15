"""
Initial Setup: Bullet indent remove - staircase effect
Task ID: osworld_impress_bullet_indent_remove_009
Domain: libreoffice_impress

Creates a 5-slide creative presentation. Slide 5 has a content textbox with 4
uniformly indented bullet items WITH bullet symbols (pre-task state).
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_bullet_indent_remove_009'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_bullet_char(para, char='•'):
    """Add a visible bullet character to the paragraph via XML."""
    pPr = para._p.get_or_add_pPr()
    # Remove any existing buNone or buChar
    for tag in (qn('a:buNone'), qn('a:buChar'), qn('a:buAutoNum')):
        existing = pPr.find(tag)
        if existing is not None:
            pPr.remove(existing)
    # Add buChar
    buChar = etree.SubElement(pPr, qn('a:buChar'))
    buChar.set('char', char)


def set_indent_level(para, margin_left_emu: int, indent_emu: int = 0):
    """
    Set paragraph left margin (marL) and hanging indent (indent) via XML.
    marL: left margin in EMU
    indent: first-line indent (negative = hanging indent)
    """
    pPr = para._p.get_or_add_pPr()
    pPr.set('marL', str(margin_left_emu))
    pPr.set('indent', str(indent_emu))


def create_initial():
    prs = Presentation()
    # Use standard widescreen (default)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Innovation in the Digital Age"
    slide1.placeholders[1].text = "A Creative Journey Through Technology & Design"
    # Subtitle formatting
    tf = slide1.placeholders[1].text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(24)
            run.font.color.rgb = RGBColor(0x44, 0x44, 0x88)

    # ---- Slide 2: Vision & Goals ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Our Vision & Goals"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Empower teams to build faster and smarter solutions"
    goals = [
        "Accelerate product development cycles by 40%",
        "Expand global reach to 25 new markets",
        "Foster a culture of continuous innovation",
        "Reduce operational overhead through automation",
    ]
    for goal in goals:
        p = tf2.add_paragraph()
        p.text = goal
        p.level = 1
        for run in p.runs:
            run.font.size = Pt(18)

    # ---- Slide 3: Key Milestones ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Key Milestones"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "2025 Roadmap Highlights"
    milestones = [
        "Q1 2025 — Platform v3.0 Launch",
        "Q2 2025 — Partnership with TechNova Corp",
        "Q3 2025 — AI Integration Rollout",
        "Q4 2025 — IPO Preparation & Investor Summit",
    ]
    for ms in milestones:
        p = tf3.add_paragraph()
        p.text = ms
        p.level = 1
        for run in p.runs:
            run.font.size = Pt(18)

    # ---- Slide 4: Team & Culture ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Team & Culture"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Our People Make the Difference"
    team_points = [
        "150+ talented engineers and designers worldwide",
        "Remote-first culture with quarterly meetups",
        "Dedicated learning budget of $5,000 per employee",
        "Diversity & Inclusion initiatives in 12 countries",
    ]
    for point in team_points:
        p = tf4.add_paragraph()
        p.text = point
        p.level = 1
        for run in p.runs:
            run.font.size = Pt(18)

    # ---- Slide 5: Creative Highlights (bullet items, uniformly indented) ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Creative Highlights"

    tf5 = slide5.placeholders[1].text_frame
    # First paragraph (existing placeholder paragraph)
    tf5.text = ""

    bullet_items = [
        "Award-winning brand identity redesign",
        "Interactive product demo experience",
        "Immersive augmented reality showcase",
        "Dynamic data visualization dashboard",
    ]

    # Uniform indentation: all items at same level with bullet symbol
    # marL = 457200 EMU (~0.5 inch = standard list indent), indent = -228600 (hanging)
    UNIFORM_MARL = 457200   # ~0.5 inch
    UNIFORM_INDENT = -228600  # hanging indent

    for i, item_text in enumerate(bullet_items):
        if i == 0:
            p = tf5.paragraphs[0]
            p.text = item_text
        else:
            p = tf5.add_paragraph()
            p.text = item_text

        for run in p.runs:
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(0x22, 0x22, 0x55)

        # Set uniform indentation
        set_indent_level(p, UNIFORM_MARL, UNIFORM_INDENT)
        # Add bullet character
        add_bullet_char(p, '•')

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup — open the file in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
