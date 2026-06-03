"""
Initial Setup: Create a 12-slide sales pitch presentation with no custom slide shows.
Task ID: impress_sales_039
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_039'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.paragraphs[0].text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
    return slide


def create_initial():
    prs = Presentation()

    # Slide 1: Title
    add_title_slide(
        prs,
        "NovaTech Solutions",
        "Full Product Pitch - Q2 2025\nPrepared by Elena Rodriguez, VP Sales"
    )

    # Slide 2: Agenda
    add_content_slide(prs, "Agenda", [
        "Market Opportunity & Industry Landscape",
        "Product Overview & Key Features",
        "Competitive Advantages",
        "Customer Success Stories",
        "Technical Architecture",
        "Pricing & Packaging",
        "Implementation Roadmap",
        "Financial Projections",
        "Partnership Opportunities",
        "Q&A"
    ])

    # Slide 3: Market Opportunity
    add_content_slide(prs, "Market Opportunity", [
        "Global SaaS market projected to reach $908B by 2030",
        "Enterprise workflow automation growing at 23.4% CAGR",
        "78% of Fortune 500 companies actively seeking AI-driven solutions",
        "Total addressable market: $45.2B in North America alone",
        "Key verticals: Healthcare, Finance, Manufacturing, Retail"
    ])

    # Slide 4: Product Overview
    add_content_slide(prs, "Product Overview - NovaTech Platform", [
        "AI-powered workflow orchestration engine",
        "200+ pre-built integrations with enterprise tools",
        "Real-time analytics dashboard with custom KPIs",
        "Role-based access control with SOC 2 compliance",
        "Multi-tenant architecture supporting 50,000+ concurrent users"
    ])

    # Slide 5: Competitive Advantages
    add_content_slide(prs, "Why NovaTech Wins", [
        "3x faster deployment vs. legacy competitors",
        "Proprietary NovaBrain AI reduces manual tasks by 67%",
        "99.99% uptime SLA - industry leading reliability",
        "30% lower TCO compared to Salesforce + ServiceNow stack",
        "Dedicated customer success team with <2hr response time"
    ])

    # Slide 6: Customer Success - Meridian Health
    add_content_slide(prs, "Customer Success: Meridian Health Systems", [
        "Challenge: 12,000 staff using 14 disconnected systems",
        "Solution: Unified NovaTech platform deployment in 8 weeks",
        "Result: 42% reduction in administrative overhead",
        "Result: $3.2M annual savings in operational costs",
        "NPS score improved from 34 to 78 within 6 months"
    ])

    # Slide 7: Technical Architecture
    add_content_slide(prs, "Technical Architecture", [
        "Microservices on Kubernetes with auto-scaling",
        "Event-driven architecture using Apache Kafka",
        "GraphQL API layer for flexible data access",
        "End-to-end encryption with AES-256 at rest",
        "Multi-region deployment: US-East, EU-West, APAC"
    ])

    # Slide 8: Pricing & Packaging
    add_content_slide(prs, "Pricing & Packaging", [
        "Starter: $29/user/month - Up to 50 users",
        "Professional: $59/user/month - Up to 500 users",
        "Enterprise: $99/user/month - Unlimited users",
        "Volume discounts: 15% for annual commitment",
        "Custom pricing available for 1,000+ seat deals"
    ])

    # Slide 9: Implementation Roadmap
    add_content_slide(prs, "Implementation Roadmap", [
        "Week 1-2: Discovery & requirements gathering",
        "Week 3-4: Platform configuration & data migration",
        "Week 5-6: Integration setup & UAT testing",
        "Week 7-8: Training rollout & go-live support",
        "Ongoing: Quarterly business reviews & optimization"
    ])

    # Slide 10: Financial Projections
    add_content_slide(prs, "Financial Projections - 3 Year Outlook", [
        "Year 1 ARR target: $12.5M (current pipeline: $18.3M)",
        "Year 2 projected ARR: $28.7M with 85% net retention",
        "Year 3 projected ARR: $52.1M targeting Series C milestone",
        "Gross margin: 82% with improving unit economics",
        "Customer acquisition cost payback: 14 months"
    ])

    # Slide 11: Partnership Opportunities
    add_content_slide(prs, "Strategic Partnership Opportunities", [
        "Technology partnerships: AWS, Azure, Google Cloud",
        "Channel partnerships: Accenture, Deloitte, PwC",
        "OEM licensing for vertical-specific solutions",
        "Co-selling programs with 2x revenue multiplier",
        "Joint innovation labs for next-gen AI capabilities"
    ])

    # Slide 12: Closing & Q&A
    add_title_slide(
        prs,
        "Let's Build the Future Together",
        "Elena Rodriguez | elena@novatech.io | (415) 555-0192\nSchedule a technical deep-dive: novatech.io/demo"
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
