"""
Initial Setup: Set a gradient background (light to dark blue) on all slides in this presentation.
Task ID: osworld_impress_all_slides_background_005
Domain: libreoffice_impress

Creates a 6-slide academic lecture deck with white backgrounds throughout.
The agent must set a gradient background on all slides.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_all_slides_background_005'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_white_background(slide):
    """Set a white solid background on a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_text(slide, title_text, body_text=None):
    """Helper to set title and body text on a slide with placeholders."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    if body_text and len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.text = body_text


def create_initial():
    prs = Presentation()

    # Use standard widescreen dimensions (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    layout0 = prs.slide_layouts[0]  # Title Slide
    slide1 = prs.slides.add_slide(layout0)
    set_white_background(slide1)
    slide1.shapes.title.text = "Introduction to Machine Learning"
    slide1.placeholders[1].text = "Department of Computer Science\nDr. Emily Hartman | Spring 2025"

    # ---- Slide 2: Course Overview ----
    layout1 = prs.slide_layouts[1]  # Title and Content
    slide2 = prs.slides.add_slide(layout1)
    set_white_background(slide2)
    slide2.shapes.title.text = "Course Overview"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "What is Machine Learning?"
    paras = [
        "Supervised Learning: labeled training data",
        "Unsupervised Learning: finding hidden patterns",
        "Reinforcement Learning: learning through interaction",
        "Deep Learning: multi-layer neural networks",
    ]
    for para_text in paras:
        p = tf2.add_paragraph()
        p.text = para_text
        p.level = 1

    # ---- Slide 3: Supervised Learning ----
    slide3 = prs.slides.add_slide(layout1)
    set_white_background(slide3)
    slide3.shapes.title.text = "Supervised Learning"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Key Concepts"
    items3 = [
        "Training set: labeled examples (X, y)",
        "Hypothesis function: maps inputs to outputs",
        "Loss function: measures prediction error",
        "Gradient descent: optimization algorithm",
        "Overfitting vs. underfitting: bias-variance tradeoff",
    ]
    for item in items3:
        p = tf3.add_paragraph()
        p.text = item
        p.level = 1

    # ---- Slide 4: Neural Networks ----
    slide4 = prs.slides.add_slide(layout1)
    set_white_background(slide4)
    slide4.shapes.title.text = "Neural Networks Architecture"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Building Blocks of Deep Learning"
    items4 = [
        "Input layer: receives raw feature vectors",
        "Hidden layers: learn intermediate representations",
        "Activation functions: ReLU, Sigmoid, Tanh",
        "Output layer: produces final predictions",
        "Backpropagation: computes gradients efficiently",
        "Dropout: regularization to prevent overfitting",
    ]
    for item in items4:
        p = tf4.add_paragraph()
        p.text = item
        p.level = 1

    # ---- Slide 5: Evaluation Metrics ----
    slide5 = prs.slides.add_slide(layout1)
    set_white_background(slide5)
    slide5.shapes.title.text = "Evaluation Metrics"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Measuring Model Performance"
    items5 = [
        "Accuracy: fraction of correct predictions",
        "Precision: TP / (TP + FP)",
        "Recall: TP / (TP + FN)",
        "F1-Score: harmonic mean of precision and recall",
        "ROC-AUC: area under the receiver operating curve",
        "Cross-validation: k-fold for robust estimation",
    ]
    for item in items5:
        p = tf5.add_paragraph()
        p.text = item
        p.level = 1

    # ---- Slide 6: Summary & Next Steps ----
    layout2 = prs.slide_layouts[2]  # Section Header or similar
    slide6 = prs.slides.add_slide(layout1)
    set_white_background(slide6)
    slide6.shapes.title.text = "Summary & Next Steps"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Key Takeaways"
    items6 = [
        "ML learns patterns from data automatically",
        "Choice of algorithm depends on problem type",
        "Feature engineering is crucial for performance",
        "Next lecture: Convolutional Neural Networks (CNNs)",
        "Assignment 2 due: March 20, 2025",
    ]
    for item in items6:
        p = tf6.add_paragraph()
        p.text = item
        p.level = 1

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open file in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
