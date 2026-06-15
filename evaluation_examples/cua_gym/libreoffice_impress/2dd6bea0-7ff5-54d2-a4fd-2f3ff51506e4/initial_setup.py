"""
Initial Setup: Create presentation with 5 animations on slide 4 set to 'With Previous' (0 delay)
Task ID: impress_fix_061
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import shutil
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_061'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_presentation():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Staggered Reveal Presentation"
    slide1.placeholders[1].text = "Q2 2025 Product Launch Strategy"

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Product Launch Overview"
    body2 = slide2.placeholders[1]
    tf2 = body2.text_frame
    tf2.text = "Our Q2 launch plan focuses on three key markets"
    for item, lvl in [
        ("North America: Enterprise SaaS platform expansion", 1),
        ("Europe: GDPR-compliant data analytics suite", 1),
        ("Asia Pacific: Mobile-first payment integration", 1),
    ]:
        p = tf2.add_paragraph()
        p.text = item
        p.level = lvl

    # --- Slide 3: Timeline ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Launch Timeline"
    body3 = slide3.placeholders[1]
    tf3 = body3.text_frame
    tf3.text = "Phase 1 (April): Beta testing with 50 enterprise clients"
    for item in [
        "Phase 2 (May): Public beta release and marketing campaign",
        "Phase 3 (June): Full production launch",
        "Phase 4 (July): Post-launch optimization and scaling",
    ]:
        p = tf3.add_paragraph()
        p.text = item

    # --- Slide 4: Key Features (animated slide) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    title_box = slide4.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.8))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "Key Features"
    p_title.alignment = PP_ALIGN.CENTER
    run_title = p_title.runs[0]
    run_title.font.size = Pt(32)
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # 5 feature text boxes that will be animated
    features = [
        ("Real-Time Analytics Dashboard", "Monitor KPIs with live data streaming and customizable widgets"),
        ("AI-Powered Recommendations", "Machine learning engine suggests actions based on user behavior patterns"),
        ("Multi-Cloud Deployment", "Deploy seamlessly across AWS, Azure, and Google Cloud with one click"),
        ("Enterprise Security Suite", "SOC 2 Type II compliant with end-to-end encryption and SSO support"),
        ("Collaborative Workspace", "Real-time co-editing with version history and conflict resolution"),
    ]

    for i, (title, desc) in enumerate(features):
        top = Inches(1.3 + i * 1.2)
        box = slide4.shapes.add_textbox(Inches(1.5), top, Inches(7), Inches(1.0))
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        r1 = p1.runs[0]
        r1.font.size = Pt(20)
        r1.font.bold = True
        r1.font.color.rgb = RGBColor(0x2E, 0x74, 0xB5)
        p2 = tf.add_paragraph()
        p2.text = desc
        r2 = p2.runs[0]
        r2.font.size = Pt(14)
        r2.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # --- Slide 5: Summary ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Next Steps"
    body5 = slide5.placeholders[1]
    tf5 = body5.text_frame
    tf5.text = "Schedule individual team demos this week"
    for item in [
        "Review pricing tiers and select enterprise package",
        "Begin integration testing with existing infrastructure",
        "Contact sales team for custom deployment options",
    ]:
        p = tf5.add_paragraph()
        p.text = item

    prs.save(OUTPUT)
    print(f'Presentation saved: {OUTPUT}')


def inject_animations():
    """Inject Fade In entrance animations on slide 4 with 'With Previous' timing and 0 delay."""
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

    slide_xml_name = 'ppt/slides/slide4.xml'

    with zipfile.ZipFile(OUTPUT, 'r') as zf:
        slide_xml = zf.read(slide_xml_name)

    root = ET.fromstring(slide_xml)
    ET.register_namespace('', ns_p)
    ET.register_namespace('a', ns_a)
    ET.register_namespace('r', ns_r)

    # Get shape IDs from slide 4
    cSld = root.find(f'{{{ns_p}}}cSld')
    spTree = cSld.find(f'{{{ns_p}}}spTree')
    shapes = spTree.findall(f'{{{ns_p}}}sp')

    shape_sp_ids = []
    for sp in shapes:
        nvSpPr = sp.find(f'{{{ns_p}}}nvSpPr')
        cNvPr = nvSpPr.find(f'{{{ns_p}}}cNvPr')
        shape_sp_ids.append(cNvPr.get('id'))

    # Last 5 shapes are the feature text boxes
    feature_ids = shape_sp_ids[-5:]
    print(f'Feature shape IDs for animation: {feature_ids}')

    # Build animation nodes: first is clickEffect (On Click), rest are withEffect (With Previous)
    # All with delay=0 so they start simultaneously
    # presetID=10 = Fade
    anim_children = []
    ctn_id = 3
    for i, sp_id in enumerate(feature_ids):
        node_type = 'clickEffect' if i == 0 else 'withEffect'
        anim_children.append(
            f'<p:par xmlns:p="{ns_p}">'
            f'<p:cTn id="{ctn_id}" fill="hold">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst><p:par>'
            f'<p:cTn id="{ctn_id+1}" fill="hold">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst><p:par>'
            f'<p:cTn id="{ctn_id+2}" presetID="10" presetClass="entr" presetSubtype="0" fill="hold" nodeType="{node_type}">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst>'
            f'<p:set><p:cBhvr>'
            f'<p:cTn id="{ctn_id+3}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
            f'<p:tgtEl><p:spTgt spid="{sp_id}"/></p:tgtEl>'
            f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
            f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
            f'<p:animEffect transition="in" filter="fade"><p:cBhvr>'
            f'<p:cTn id="{ctn_id+4}" dur="500"/>'
            f'<p:tgtEl><p:spTgt spid="{sp_id}"/></p:tgtEl>'
            f'</p:cBhvr></p:animEffect>'
            f'</p:childTnLst></p:cTn>'
            f'</p:par></p:childTnLst></p:cTn>'
            f'</p:par></p:childTnLst></p:cTn>'
            f'</p:par>'
        )
        ctn_id += 5

    timing_str = (
        f'<p:timing xmlns:p="{ns_p}">'
        f'<p:tnLst><p:par>'
        f'<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
        f'<p:childTnLst>'
        f'<p:seq concurrent="1" nextAc="seek">'
        f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq">'
        f'<p:childTnLst>{"".join(anim_children)}</p:childTnLst>'
        f'</p:cTn>'
        f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
        f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
        f'</p:seq>'
        f'</p:childTnLst></p:cTn>'
        f'</p:par></p:tnLst>'
        f'</p:timing>'
    )

    timing_el = ET.fromstring(timing_str)
    root.append(timing_el)

    new_xml = '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n' + ET.tostring(root, encoding='unicode')

    # Rewrite the zip
    tmp_path = OUTPUT + '.tmp'
    shutil.copy(OUTPUT, tmp_path)
    with zipfile.ZipFile(tmp_path, 'r') as zin:
        with zipfile.ZipFile(OUTPUT, 'w') as zout:
            for item in zin.infolist():
                if item.filename == slide_xml_name:
                    zout.writestr(item, new_xml)
                else:
                    zout.writestr(item, zin.read(item.filename))
    os.remove(tmp_path)
    print(f'Animations injected into slide 4 (5x Fade In, With Previous, delay=0)')


def main():
    create_presentation()
    inject_animations()
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


main()
