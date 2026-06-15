"""
Reward Script: Investor Update Presentation
Task ID: impress_wf_049
Domain: libreoffice_impress
Scoring:
  C1: 10 slides (0.10)
  C2: Slide 1 title text (0.10)
  C3: Background color #E8EAF6 (0.10)
  C4: Slide 2 metric cards with triangles (0.10)
  C5: Slide 3 LINE chart (0.15)
  C6: Slide 4 DOUGHNUT chart (0.15)
  C7: Slide 5 churn analysis shapes (0.10)
  C8: Slide 7 quadrant lines (0.05)
  C9: Slide 8 financial summary table (0.15)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_049'
FILE_NAME = 'Investor_Update.pptx'


def get_slide_bg_rgb(slide):
    """Get background RGB color of a slide, handling inheritance."""
    try:
        fill = slide.background.fill
        if fill.type == 1:  # solid
            return str(fill.fore_color.rgb)
        elif fill.type == 5:  # inherited from master
            master_fill = slide.slide_layout.slide_master.background.fill
            if master_fill.type == 1:
                return str(master_fill.fore_color.rgb)
    except Exception:
        pass
    return None


def get_all_text(slide):
    """Extract all text from a slide including grouped shapes."""
    texts = []
    def extract(shape):
        if hasattr(shape, 'text') and hasattr(shape, 'text_frame'):
            texts.append(shape.text)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                extract(sub)
    for shape in slide.shapes:
        extract(shape)
    return texts


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)

    # Component 1: Exactly 10 slides (0.10 points)
    try:
        slide_count = len(slides)
        if slide_count == 10:
            print(f"PASS: Component 1 — 10 slides found (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 1 — expected 10 slides, found {slide_count}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 1 title contains 'Q3 2024 Investor Update' and 'CloudBase' (0.10 points)
    try:
        if len(slides) >= 1:
            slide1_texts = get_all_text(slides[0])
            full_text = " ".join(slide1_texts).lower()
            has_title = "q3 2024 investor update" in full_text
            has_company = "cloudbase" in full_text
            if has_title and has_company:
                print(f"PASS: Component 2 — Slide 1 has correct title and company name (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 2 — title={'found' if has_title else 'missing'}, company={'found' if has_company else 'missing'}")
        else:
            print(f"FAIL: Component 2 — no slides")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Background color #E8EAF6 on majority of slides (0.10 points)
    try:
        bg_match_count = 0
        for slide in slides:
            bg = get_slide_bg_rgb(slide)
            if bg and bg.upper() == "E8EAF6":
                bg_match_count += 1
        # At least 8 of 10 slides should have the correct background
        if bg_match_count >= 8:
            print(f"PASS: Component 3 — {bg_match_count}/{len(slides)} slides have #E8EAF6 background (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 3 — only {bg_match_count}/{len(slides)} slides have #E8EAF6 background")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 2 has metric cards with triangle shapes (0.10 points)
    try:
        if len(slides) >= 2:
            slide2 = slides[1]
            triangle_count = 0
            rounded_rect_count = 0
            for shape in slide2.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    name_lower = shape.name.lower()
                    if "triangle" in name_lower:
                        triangle_count += 1
                    if "rounded" in name_lower or "rectangle" in name_lower:
                        rounded_rect_count += 1

            slide2_texts = get_all_text(slide2)
            full_text = " ".join(slide2_texts).lower()
            has_highlight_title = "highlight" in full_text

            # Need at least 2 triangles (up/down indicators) and some card shapes
            if triangle_count >= 2 and (rounded_rect_count >= 2 or has_highlight_title):
                print(f"PASS: Component 4 — Slide 2 has {triangle_count} triangles and {rounded_rect_count} card rects (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 4 — triangles={triangle_count}, rects={rounded_rect_count}, highlight_title={has_highlight_title}")
        else:
            print(f"FAIL: Component 4 — not enough slides")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slide 3 has a LINE chart (MRR growth) (0.15 points)
    try:
        if len(slides) >= 3:
            slide3 = slides[2]
            has_line_chart = False
            for shape in slide3.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    # LINE chart type is 4
                    chart_type_val = shape.chart.chart_type
                    if chart_type_val == 4:  # LINE
                        has_line_chart = True

            slide3_texts = get_all_text(slide3)
            full_text = " ".join(slide3_texts).lower()
            has_mrr_text = "mrr" in full_text

            if has_line_chart and has_mrr_text:
                print(f"PASS: Component 5 — Slide 3 has LINE chart with MRR text (0.15 pts)")
                total_score += 0.15
            elif has_line_chart:
                print(f"PARTIAL: Component 5 — Slide 3 has LINE chart but missing MRR text (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 5 — no LINE chart found on slide 3")
        else:
            print(f"FAIL: Component 5 — not enough slides")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Slide 4 has a DOUGHNUT chart (customer segments) (0.15 points)
    try:
        if len(slides) >= 4:
            slide4 = slides[3]
            has_doughnut = False
            for shape in slide4.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    chart_type_val = shape.chart.chart_type
                    if chart_type_val == -4120:  # DOUGHNUT
                        has_doughnut = True

            slide4_texts = get_all_text(slide4)
            full_text = " ".join(slide4_texts).lower()
            has_segment_text = "segment" in full_text or "customer" in full_text or "breakdown" in full_text

            if has_doughnut and has_segment_text:
                print(f"PASS: Component 6 — Slide 4 has DOUGHNUT chart with segment text (0.15 pts)")
                total_score += 0.15
            elif has_doughnut:
                print(f"PARTIAL: Component 6 — Slide 4 has DOUGHNUT chart but missing segment text (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 6 — no DOUGHNUT chart found on slide 4")
        else:
            print(f"FAIL: Component 6 — not enough slides")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Slide 5 has churn analysis with shapes (bars + lines simulated) (0.10 points)
    try:
        if len(slides) >= 5:
            slide5 = slides[4]
            auto_shape_count = sum(1 for s in slide5.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE)
            line_count = sum(1 for s in slide5.shapes if s.shape_type == MSO_SHAPE_TYPE.LINE)

            slide5_texts = get_all_text(slide5)
            full_text = " ".join(slide5_texts).lower()
            has_churn_text = "churn" in full_text

            # Churn analysis uses bars (auto shapes) and lines to simulate dual-axis chart
            if auto_shape_count >= 5 and line_count >= 4 and has_churn_text:
                print(f"PASS: Component 7 — Slide 5 has churn analysis ({auto_shape_count} shapes, {line_count} lines) (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 7 — auto_shapes={auto_shape_count}, lines={line_count}, churn_text={has_churn_text}")
        else:
            print(f"FAIL: Component 7 — not enough slides")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    # Component 8: Slide 7 has competitive positioning quadrant (lines + labels) (0.05 points)
    try:
        if len(slides) >= 7:
            slide7 = slides[6]
            line_count = sum(1 for s in slide7.shapes if s.shape_type == MSO_SHAPE_TYPE.LINE)

            slide7_texts = get_all_text(slide7)
            full_text = " ".join(slide7_texts).lower()
            has_competitive = "competitive" in full_text or "positioning" in full_text
            has_quadrant_labels = ("leader" in full_text or "challenger" in full_text)

            if line_count >= 2 and has_competitive and has_quadrant_labels:
                print(f"PASS: Component 8 — Slide 7 has quadrant ({line_count} lines, labels present) (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 8 — lines={line_count}, competitive={has_competitive}, quadrant_labels={has_quadrant_labels}")
        else:
            print(f"FAIL: Component 8 — not enough slides")
    except Exception as e:
        print(f"ERROR: Component 8 — {e}")

    # Component 9: Slide 8 has financial summary table (0.15 points)
    try:
        if len(slides) >= 8:
            slide8 = slides[7]
            has_table = False
            table_rows = 0
            table_cols = 0
            has_financial_header = False

            for shape in slide8.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    has_table = True
                    t = shape.table
                    table_rows = len(t.rows)
                    table_cols = len(t.columns)
                    # Check header row has financial columns
                    header_text = " ".join(t.cell(0, c).text.lower() for c in range(len(t.columns)))
                    if "q1" in header_text or "q2" in header_text or "q3" in header_text:
                        has_financial_header = True

            slide8_texts = get_all_text(slide8)
            full_text = " ".join(slide8_texts).lower()
            has_financial_title = "financial" in full_text

            if has_table and table_rows >= 5 and table_cols >= 4 and has_financial_header:
                print(f"PASS: Component 9 — Slide 8 has financial table ({table_rows}x{table_cols}) (0.15 pts)")
                total_score += 0.15
            elif has_table:
                print(f"PARTIAL: Component 9 — Table found but size/headers insufficient ({table_rows}x{table_cols}, fin_hdr={has_financial_header}) (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 9 — no table found on slide 8")
        else:
            print(f"FAIL: Component 9 — not enough slides")
    except Exception as e:
        print(f"ERROR: Component 9 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/Desktop/{FILE_NAME}'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
