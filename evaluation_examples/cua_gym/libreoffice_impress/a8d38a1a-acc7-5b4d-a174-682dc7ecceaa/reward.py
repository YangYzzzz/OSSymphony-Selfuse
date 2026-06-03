"""
FINAL REWARD SCRIPT - SUCCESS
Task: I’m polishing up a big deck and just spotted that slide 61 still says “TBD” in a couple of spots. Using LibreOffice Impress, how do I quickly change every single instance of “TBD” on that specific slide to “Final” so nothing slips through the cracks?
Generated: 2025-09-10 22:28:32
Status: success
Model: azure-o3
Total Steps: 1
"""

from pptx import Presentation
import os, re

def verify_replace_tbd_with_final(file_path: str) -> float:
    """
    Verify that, on slide 61 (index 60) of the provided PPTX file,
    every instance of the text 'TBD' has been replaced by 'Final'.

    Scoring (progressive):
      • 0.7 points  – slide 61 contains NO occurrence of 'TBD' (case-insensitive)
      • 0.3 points  – slide 61 DOES contain at least one occurrence of 'Final'
                      *and* still no 'TBD' (ensures true replacement, not mere addition)
      = 1.0 points  – both conditions satisfied
    """

    print(f"Verifying presentation: {file_path}")
    score = 0.0

    # ---------- 1) Basic file checks (no points awarded) ----------
    if not os.path.exists(file_path):
        print("✗ File not found – cannot verify task")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as exc:
        print(f"✗ Could not open presentation: {exc}")
        return 0.0

    # ---------- 2) Ensure slide 61 exists (no points) ----------
    slide_index = 60  # zero-based index
    if len(prs.slides) <= slide_index:
        print(f"✗ Slide 61 missing – only {len(prs.slides)} slide(s) present")
        return 0.0

    slide = prs.slides[slide_index]
    print("✓ Slide 61 loaded successfully")

    # ---------- 3) Collect all text on slide 61 ----------
    slide_texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            txt = shape.text
            if txt:
                slide_texts.append(txt)
                print(f"  • Found text: {txt!r}")

    if not slide_texts:
        print("✗ No text found on slide 61 – nothing to verify")
        return 0.0

    combined_text = "\n".join(slide_texts).lower()

    # ---------- 4) Verification & progressive scoring ----------
    weight_no_tbd  = 0.7
    weight_has_final = 0.3

    # 4a) Ensure *no* 'tbd' remains
    if "tbd" not in combined_text:
        print("✓ No occurrences of 'TBD' detected on slide 61")
        score += weight_no_tbd
    else:
        remaining = len(re.findall(r"tbd", combined_text))
        print(f"✗ {remaining} occurrence(s) of 'TBD' still present")

    # 4b) Ensure 'final' appears somewhere (and still no 'tbd')
    if "final" in combined_text and "tbd" not in combined_text:
        print("✓ 'Final' detected on slide 61")
        score += weight_has_final
    elif "final" not in combined_text:
        print("✗ No occurrence of 'Final' found on slide 61")

    final_score = min(score, 1.0)
    print(f"Total score: {final_score}/1.0")
    return final_score


if __name__ == "__main__":
    pptx_path = "/home/user/im_polishing_up_a_big_deck_and_just_spotted_that_slide_61_still_says_tbd_in_a_couple_of_spots_using__golden.pptx"
    reward = verify_replace_tbd_with_final(pptx_path)
    print(f"REWARD: {reward}")
