"""
Initial Setup: Build a concept map on slide 4 of Plant_Biology.pptx
Task ID: impress_teach_053
Domain: libreoffice_impress

Creates a 7-slide Plant Biology presentation. Slide 4 has only the title
'Photosynthesis Overview' and is otherwise empty (the agent must add the concept map).
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_053'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(18)
    return slide


def add_title_only_slide(prs, title_text):
    """Add a slide with only a title placeholder and no body content."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add a text box as the title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(prs, "Plant Biology", "An Introduction to Botanical Sciences\nDr. Elena Ramirez | Spring 2025")

    # Slide 2: Cell Structure
    add_content_slide(prs, "Plant Cell Structure", [
        "Cell wall provides rigid structural support made of cellulose",
        "Chloroplasts contain chlorophyll for photosynthesis",
        "Large central vacuole maintains turgor pressure",
        "Plasmodesmata enable cell-to-cell communication",
        "Endoplasmic reticulum synthesizes proteins and lipids",
        "Mitochondria generate ATP through cellular respiration",
    ])

    # Slide 3: Plant Reproduction
    add_content_slide(prs, "Plant Reproduction", [
        "Sexual reproduction involves pollination and seed formation",
        "Asexual reproduction through vegetative propagation",
        "Flowers attract pollinators with color and nectar",
        "Seeds contain embryo, endosperm, and seed coat",
        "Germination requires water, oxygen, and temperature",
    ])

    # Slide 4: Photosynthesis Overview - EMPTY (only title)
    add_title_only_slide(prs, "Photosynthesis Overview")

    # Slide 5: Nutrient Transport
    add_content_slide(prs, "Nutrient Transport in Plants", [
        "Xylem transports water and minerals upward from roots",
        "Phloem distributes sugars from leaves to all plant parts",
        "Root hairs increase surface area for water absorption",
        "Transpiration creates negative pressure pulling water up",
        "Guard cells regulate stomatal opening and gas exchange",
    ])

    # Slide 6: Plant Hormones
    add_content_slide(prs, "Plant Hormones and Growth", [
        "Auxin promotes cell elongation and phototropism",
        "Gibberellins stimulate stem elongation and seed germination",
        "Cytokinins promote cell division and delay senescence",
        "Ethylene triggers fruit ripening and leaf abscission",
        "Abscisic acid regulates dormancy and stomatal closure",
    ])

    # Slide 7: Summary
    add_content_slide(prs, "Summary and Key Takeaways", [
        "Plants have specialized cells with unique organelles",
        "Photosynthesis converts light energy to chemical energy",
        "Vascular tissue enables long-distance transport",
        "Hormones coordinate growth, development, and responses",
        "Understanding plant biology is crucial for agriculture",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
