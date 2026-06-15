"""
Reward Script: Build a concept map on slide 4 with ovals and connector arrows
Task ID: impress_teach_053
Domain: libreoffice_impress
Scoring:
  Component 1: Central oval 'Photosynthesis' with fill #81C784 (0.25)
  Component 2: Four surrounding ovals with correct labels and fill #E8F5E9 (0.40)
  Component 3: All concept-map ovals use ellipse geometry (0.10)
  Component 4: Four connector arrows with arrowhead markers (0.25)
"""

import os
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_053'


def persist_app_state(domain):
    """Save any unsaved GUI state before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_shape_geometry(shape):
    """Return the preset geometry name (e.g. 'ellipse', 'rect') or None."""
    try:
        sp_pr = shape.element.find(qn('p:spPr'))
        if sp_pr is not None:
            prst_geom = sp_pr.find(qn('a:prstGeom'))
            if prst_geom is not None:
                return prst_geom.get('prst')
    except Exception:
        pass
    return None


def get_shape_fill_rgb(shape):
    """Return fill RGB as uppercase hex string (e.g. '81C784') or None."""
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb).upper()
    except Exception:
        pass
    return None


def get_shape_text(shape):
    """Return combined text from all paragraphs, stripped."""
    if not shape.has_text_frame:
        return ""
    parts = []
    for para in shape.text_frame.paragraphs:
        parts.append(para.text)
    return " ".join(parts).strip()


def has_arrow_marker(shape):
    """Check if a connector/line shape has an arrow marker (headEnd or tailEnd with type='arrow')."""
    try:
        xml_str = shape.element.xml if hasattr(shape.element, 'xml') else ""
        # Fall back to searching child elements
        for ln in shape.element.findall('.//' + qn('a:ln')):
            for child in ln:
                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                if tag in ('headEnd', 'tailEnd'):
                    arrow_type = child.get('type', '')
                    if arrow_type and arrow_type != 'none':
                        return True
    except Exception:
        pass
    return False


def verify_task(file_path):
    """
    Verify concept map on slide 4 with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # slide 4 (0-indexed)

    # Collect auto shapes (potential ovals) and connectors from slide 4
    auto_shapes = []
    connectors = []
    for shape in slide.shapes:
        if shape.shape_type == 1:  # AUTO_SHAPE
            auto_shapes.append(shape)
        elif shape.shape_type == 9:  # LINE/connector
            connectors.append(shape)

    # ------------------------------------------------------------------
    # Component 1: Central oval 'Photosynthesis' with fill #81C784 (0.25)
    # ------------------------------------------------------------------
    try:
        center_match = [s for s in auto_shapes
                        if get_shape_text(s).strip().lower() == 'photosynthesis'
                        and get_shape_fill_rgb(s) == '81C784']
        if len(center_match) > 0:
            print(f"PASS: Component 1 -- Central oval 'Photosynthesis' with fill #81C784 found (0.25 pts)")
            total_score += 0.25
        else:
            # Check if text exists but wrong fill
            text_match = [s for s in auto_shapes if get_shape_text(s).strip().lower() == 'photosynthesis']
            if text_match:
                fill_rgb = get_shape_fill_rgb(text_match[0])
                print(f"FAIL: Component 1 -- Found 'Photosynthesis' oval but fill is {fill_rgb}, expected 81C784")
            else:
                print(f"FAIL: Component 1 -- No oval with text 'Photosynthesis' and fill #81C784 found")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # ------------------------------------------------------------------
    # Component 2: Four surrounding ovals with correct labels and fill #E8F5E9 (0.40)
    # Each correct outer oval = 0.10 points
    # ------------------------------------------------------------------
    try:
        expected_labels = {'sunlight', 'water', 'co2', 'glucose'}
        found_labels = set()
        for shape in auto_shapes:
            text = get_shape_text(shape).strip().lower()
            if text in expected_labels:
                fill_rgb = get_shape_fill_rgb(shape)
                if fill_rgb == 'E8F5E9':
                    found_labels.add(text)

        if found_labels == expected_labels:
            print(f"PASS: Component 2 -- All 4 outer ovals found with correct labels and fill #E8F5E9 (0.40 pts)")
            total_score += 0.40
        elif len(found_labels) > 0:
            partial_pts = len(found_labels) * 0.10
            print(f"PARTIAL: Component 2 -- {len(found_labels)}/4 outer ovals correct: {found_labels} ({partial_pts} pts)")
            total_score += partial_pts
        else:
            print(f"FAIL: Component 2 -- No outer ovals with expected labels and fill #E8F5E9 found")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # ------------------------------------------------------------------
    # Component 3: All concept-map ovals use ellipse geometry (0.10)
    # ------------------------------------------------------------------
    try:
        all_labels = {'photosynthesis', 'sunlight', 'water', 'co2', 'glucose'}
        oval_shapes = [s for s in auto_shapes if get_shape_text(s).strip().lower() in all_labels]
        if len(oval_shapes) >= 5:
            all_ellipse = all(get_shape_geometry(s) == 'ellipse' for s in oval_shapes)
            if all_ellipse:
                print(f"PASS: Component 3 -- All {len(oval_shapes)} concept ovals are ellipse geometry (0.10 pts)")
                total_score += 0.10
            else:
                geoms = [get_shape_geometry(s) for s in oval_shapes]
                print(f"FAIL: Component 3 -- Not all ovals are ellipse. Geometries: {geoms}")
        else:
            print(f"FAIL: Component 3 -- Only {len(oval_shapes)} concept ovals found, need 5 for geometry check")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # ------------------------------------------------------------------
    # Component 4: Four connector arrows with arrowhead markers (0.25)
    # Each arrow connector = 0.0625 points
    # ------------------------------------------------------------------
    try:
        arrow_connectors = [c for c in connectors if has_arrow_marker(c)]
        if len(arrow_connectors) >= 4:
            print(f"PASS: Component 4 -- {len(arrow_connectors)} arrow connectors found (>= 4 required) (0.25 pts)")
            total_score += 0.25
        elif len(arrow_connectors) > 0:
            partial = len(arrow_connectors) * 0.0625
            print(f"PARTIAL: Component 4 -- {len(arrow_connectors)}/4 arrow connectors found ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 -- No arrow connectors found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
