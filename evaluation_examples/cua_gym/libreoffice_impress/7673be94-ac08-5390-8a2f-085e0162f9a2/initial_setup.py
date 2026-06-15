"""
Initial Setup: Annual Report Presentation (6 slides, mixed photo/text)
Task ID: osworld_impress_conditional_bg_image_006
Domain: libreoffice_impress

Creates a 6-slide annual report presentation:
- Slides 1, 3, 5: photo slides (white background, simulated photo content)
- Slides 2, 4, 6: text-only slides (white background)
- Slide 1 title: 'Annual Report'
- All backgrounds: white
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_conditional_bg_image_006'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_white_background(slide):
    """Set a slide's background to solid white."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_photo_placeholder(slide, slide_num, prs):
    """Add simulated photo content (colored rectangle + caption) to a photo slide."""
    # Simulated photo: a colored rectangle filling most of the slide
    photo_colors = [
        RGBColor(0xB0, 0xC4, 0xDE),   # slide 1 - light blue
        RGBColor(0x8F, 0xBC, 0x8F),   # slide 3 - dark sea green
        RGBColor(0xDC, 0xA8, 0x78),   # slide 5 - sandy brown
    ]
    idx = [1, 3, 5].index(slide_num)
    color = photo_colors[idx]

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Add rectangle shape representing a photo
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.5), Inches(1.6),
        Inches(9.0), Inches(4.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    # Photo label inside the rectangle
    photo_labels = {
        1: "Q1 Office Expansion — Kuala Lumpur HQ",
        3: "Product Launch Event — March 2024",
        5: "Team Retreat — Penang, November 2024",
    }
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = photo_labels[slide_num]
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    run.font.bold = True


def create_initial():
    prs = Presentation()

    # Use standard widescreen: 10x7.5 inches
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Photo Slide — Title + photo content ----
    layout_title = prs.slide_layouts[0]  # Title Slide layout
    slide1 = prs.slides.add_slide(layout_title)
    set_white_background(slide1)

    # Title placeholder
    title_ph = slide1.shapes.title
    title_ph.text = "Annual Report"
    for para in title_ph.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    # Subtitle placeholder
    if len(slide1.placeholders) > 1:
        sub_ph = slide1.placeholders[1]
        sub_ph.text = "Fiscal Year 2024"
        for para in sub_ph.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(22)
                run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    add_photo_placeholder(slide1, 1, prs)

    # ---- Slide 2: Text Slide — Executive Summary ----
    layout_content = prs.slide_layouts[1]  # Title + Content layout
    slide2 = prs.slides.add_slide(layout_content)
    set_white_background(slide2)

    slide2.shapes.title.text = "Executive Summary"
    for para in slide2.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    content_ph = slide2.placeholders[1]
    content_tf = content_ph.text_frame
    content_tf.text = "Revenue grew 18% year-over-year to $4.2B"
    bullets = [
        "Operating income increased to $820M (+22%)",
        "Headcount expanded to 12,400 globally",
        "Launched 3 new product lines in APAC",
        "Achieved carbon-neutral operations in Q3",
        "Returned $350M to shareholders via buybacks",
    ]
    for b in bullets:
        p = content_tf.add_paragraph()
        p.text = b
        p.level = 1
        for run in p.runs:
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ---- Slide 3: Photo Slide — Product Highlights ----
    slide3 = prs.slides.add_slide(layout_content)
    set_white_background(slide3)

    slide3.shapes.title.text = "Product Highlights"
    for para in slide3.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    add_photo_placeholder(slide3, 3, prs)

    # ---- Slide 4: Text Slide — Financial Performance ----
    slide4 = prs.slides.add_slide(layout_content)
    set_white_background(slide4)

    slide4.shapes.title.text = "Financial Performance"
    for para in slide4.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    fin_ph = slide4.placeholders[1]
    fin_tf = fin_ph.text_frame
    fin_tf.text = "Full Year 2024 Financial Highlights"
    fin_data = [
        "Total Revenue: $4.24B  (FY2023: $3.59B)",
        "Gross Margin: 61.3%  (FY2023: 58.7%)",
        "EBITDA: $1.02B  (FY2023: $0.84B)",
        "Free Cash Flow: $740M  (FY2023: $590M)",
        "EPS (diluted): $6.82  (FY2023: $5.51)",
    ]
    for fd in fin_data:
        p = fin_tf.add_paragraph()
        p.text = fd
        p.level = 1
        for run in p.runs:
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ---- Slide 5: Photo Slide — People & Culture ----
    slide5 = prs.slides.add_slide(layout_content)
    set_white_background(slide5)

    slide5.shapes.title.text = "People & Culture"
    for para in slide5.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    add_photo_placeholder(slide5, 5, prs)

    # ---- Slide 6: Text Slide — Outlook 2025 ----
    slide6 = prs.slides.add_slide(layout_content)
    set_white_background(slide6)

    slide6.shapes.title.text = "Outlook 2025"
    for para in slide6.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    out_ph = slide6.placeholders[1]
    out_tf = out_ph.text_frame
    out_tf.text = "Strategic Priorities for FY2025"
    outlook = [
        "Expand APAC and MENA market presence",
        "Invest $500M in R&D for AI-driven products",
        "Scale headcount to 15,000 by end of year",
        "Target revenue of $5.1B with 63% gross margin",
        "Complete acquisition of DataBridge Technologies",
    ]
    for ot in outlook:
        p = out_tf.add_paragraph()
        p.text = ot
        p.level = 1
        for run in p.runs:
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
