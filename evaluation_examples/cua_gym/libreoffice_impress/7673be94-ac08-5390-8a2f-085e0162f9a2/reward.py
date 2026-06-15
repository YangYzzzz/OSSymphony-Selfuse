"""
Reward Script: Apply amber background to photo slides and update title
Task ID: osworld_impress_conditional_bg_image_006
Domain: libreoffice_impress
Scoring:
  - Component 1: Slide 1 has amber (#FFBF00) background (0.2 pts)
  - Component 2: Slide 3 has amber (#FFBF00) background (0.2 pts)
  - Component 3: Slide 5 has amber (#FFBF00) background (0.2 pts)
  - Component 4: Photo slides amber AND text slides white (compound check, ensures no spillover) (0.1 pts)
  - Component 5: Slide 1 title updated to 'Visual Annual Report 2024' (0.3 pts)
  Total: 1.0
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_conditional_bg_image_006'

TARGET_BG_COLOR = 'FFBF00'  # warm amber
WHITE_BG_COLOR = 'FFFFFF'
EXPECTED_TITLE = 'Visual Annual Report 2024'


def get_slide_background_rgb(slide):
    """Return the RGBColor of slide background, or None if not solid."""
    fill = slide.background.fill
    if fill.type == 1:  # SOLID
        try:
            return fill.fore_color.rgb
        except Exception:
            return None
    elif fill.type == 5:  # inherited from master
        try:
            master_fill = slide.slide_layout.slide_master.background.fill
            if master_fill.type == 1:
                return master_fill.fore_color.rgb
        except Exception:
            return None
    return None


def get_slide_title_text(slide):
    """Return the title text from a slide, checking all text shapes."""
    # Try placeholder title first
    for shape in slide.shapes:
        if shape.has_text_frame and shape.shape_type == 13:
            pass  # skip pictures
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.idx == 0:  # title placeholder
                return shape.text_frame.text.strip()
    # Fallback: return first non-empty text shape
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                return text
    return ""


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: ensure 6 slides exist
    if len(prs.slides) != 6:
        print(f"CRITICAL: Expected 6 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide 1 (index 0) has amber (#FFBF00) background (0.2 pts)
    try:
        slide1_rgb = get_slide_background_rgb(prs.slides[0])
        slide1_rgb_str = str(slide1_rgb) if slide1_rgb is not None else None
        if slide1_rgb_str == TARGET_BG_COLOR:
            print(f"PASS: Component 1 — Slide 1 background is #{TARGET_BG_COLOR} (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 1 — Slide 1 background expected #{TARGET_BG_COLOR}, found #{slide1_rgb_str}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 3 (index 2) has amber (#FFBF00) background (0.2 pts)
    try:
        slide3_rgb = get_slide_background_rgb(prs.slides[2])
        slide3_rgb_str = str(slide3_rgb) if slide3_rgb is not None else None
        if slide3_rgb_str == TARGET_BG_COLOR:
            print(f"PASS: Component 2 — Slide 3 background is #{TARGET_BG_COLOR} (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 2 — Slide 3 background expected #{TARGET_BG_COLOR}, found #{slide3_rgb_str}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 5 (index 4) has amber (#FFBF00) background (0.2 pts)
    try:
        slide5_rgb = get_slide_background_rgb(prs.slides[4])
        slide5_rgb_str = str(slide5_rgb) if slide5_rgb is not None else None
        if slide5_rgb_str == TARGET_BG_COLOR:
            print(f"PASS: Component 3 — Slide 5 background is #{TARGET_BG_COLOR} (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 3 — Slide 5 background expected #{TARGET_BG_COLOR}, found #{slide5_rgb_str}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: All three photo slides (1, 3, 5) have amber background AND text slides (2, 4, 6) are white (0.1 pts)
    # This is a compound check: both conditions must hold — the photo slides changed color AND text slides remain white
    # This FAILS on initial_env (photo slides are still white there) and PASSES on golden_env
    try:
        photo_slide_indices = [0, 2, 4]  # slides 1, 3, 5 (0-indexed)
        text_slide_indices = [1, 3, 5]   # slides 2, 4, 6 (0-indexed)

        # All photo slides must be amber
        photo_slides_amber = all(
            str(get_slide_background_rgb(prs.slides[idx])) == TARGET_BG_COLOR
            for idx in photo_slide_indices
        )
        # All text slides must be white
        text_slides_white = all(
            str(get_slide_background_rgb(prs.slides[idx])) == WHITE_BG_COLOR
            for idx in text_slide_indices
        )

        if photo_slides_amber and text_slides_white:
            print(f"PASS: Component 4 — Photo slides amber AND text slides white (0.1 pts)")
            total_score += 0.1
        elif not photo_slides_amber:
            print(f"FAIL: Component 4 — Not all photo slides (1,3,5) have amber background yet")
        else:
            print(f"FAIL: Component 4 — Some text-only slides (2,4,6) have non-white background")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slide 1 title updated to 'Visual Annual Report 2024' (0.3 pts)
    try:
        title_text = get_slide_title_text(prs.slides[0])
        if title_text == EXPECTED_TITLE:
            print(f"PASS: Component 5 — Slide 1 title is '{title_text}' (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 5 — Slide 1 title expected '{EXPECTED_TITLE}', found '{title_text}'")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
