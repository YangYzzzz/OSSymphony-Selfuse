"""
Reward Script: Insert English and French notes into French_Lesson.pptx slides
Task ID: osworld_multi_apps_impress_notes_import_013
Domain: libreoffice_impress
Scoring:
  Component 1: All 7 slides have non-empty notes text (0.2 pts)
  Component 2: All 7 slides contain '---' separator in notes (0.3 pts)
  Component 3: English note text correctly placed (first part before '---') in >= 5/7 slides (0.25 pts)
  Component 4: French note text correctly placed (part after '---') in >= 5/7 slides (0.25 pts)
  Total: 1.0
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'osworld_multi_apps_impress_notes_import_013'
PPTX_PATH = '/home/user/Desktop/French_Lesson.pptx'
NOTES_EN_PATH = '/home/user/Desktop/french_notes_en.docx'
NOTES_FR_PATH = '/home/user/Desktop/french_notes_fr.docx'

# Expected notes per slide (English), index 0 = slide 1
ENGLISH_NOTES = [
    "Welcome students to the first French lesson. Explain the course structure: 7 sessions covering greetings, numbers, colors, days, and classroom language. Ask students about their motivation for learning French.",
    "Drill each greeting with the class. Emphasize that 'Bonjour' is used until approximately 6 PM, after which 'Bonsoir' is more appropriate. Note that 'Salut' is informal and should not be used with strangers or teachers.",
    "Have each student introduce themselves using the phrases on screen. Remind them that 'Enchanté' is masculine and 'Enchantée' is feminine. Practice: Je m'appelle [Name], j'ai [age] ans, je suis de [city].",
    "Count together as a class from 1 to 10 and back. Point out the pronunciation differences: 'cinq' is /sɛ̃k/, 'six' is /sis/ in isolation but /si/ before a consonant. Use flashcards for a quick memory game.",
    "Show color swatches alongside the French words. Note that adjectives in French agree in gender and number with the noun. Example: 'un ballon rouge' (masculine) vs 'une voiture rouge' (feminine). Colors ending in 'e' don't change form.",
    "Note that the French week starts on Monday, not Sunday. Lundi through Vendredi are weekdays; Samedi and Dimanche are the weekend. Ask students to name their favorite day in French and explain why.",
    "Practice each classroom phrase aloud. Encourage students to use these expressions throughout the rest of the course. Stress that asking for clarification is a sign of engagement, not weakness. End the lesson with a brief review quiz."
]

# Expected notes per slide (French), index 0 = slide 1
FRENCH_NOTES = [
    "Bienvenue aux étudiants dans ce premier cours de français. Expliquez la structure du cours : 7 séances couvrant les salutations, les chiffres, les couleurs, les jours et le langage de classe. Demandez aux étudiants leur motivation pour apprendre le français.",
    "Pratiquez chaque salutation avec la classe. Soulignez que 'Bonjour' s'utilise jusqu'à environ 18h, après quoi 'Bonsoir' est plus approprié. Notez que 'Salut' est informel et ne doit pas être utilisé avec des inconnus ou des professeurs.",
    "Demandez à chaque étudiant de se présenter en utilisant les expressions à l'écran. Rappelez-leur qu''Enchanté' est masculin et 'Enchantée' est féminin. Pratique : Je m'appelle [Prénom], j'ai [âge] ans, je suis de [ville].",
    "Comptez ensemble de 1 à 10 et retour. Signalez les différences de prononciation : 'cinq' se prononce /sɛ̃k/, 'six' se prononce /sis/ seul mais /si/ devant une consonne. Utilisez des cartes mémoire pour un jeu rapide de mémorisation.",
    "Montrez des échantillons de couleurs avec les mots français. Notez que les adjectifs en français s'accordent en genre et en nombre avec le nom. Exemple : 'un ballon rouge' (masculin) vs 'une voiture rouge' (féminin). Les couleurs se terminant par 'e' ne changent pas de forme.",
    "Notez que la semaine française commence le lundi, pas le dimanche. Du lundi au vendredi, ce sont les jours de semaine ; samedi et dimanche sont le week-end. Demandez aux étudiants de nommer leur jour préféré en français et d'expliquer pourquoi.",
    "Pratiquez chaque expression de classe à voix haute. Encouragez les étudiants à utiliser ces expressions tout au long du reste du cours. Insistez sur le fait que demander des éclaircissements est un signe d'engagement, pas de faiblesse. Terminez la leçon par un quiz de révision rapide."
]


def verify_task(pptx_path):
    """
    Verify that French_Lesson.pptx has notes imported from both docx files.
    Notes format: [English text] + newline + '---' + newline + [French text]
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: file must exist and be loadable
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load presentation {pptx_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must be 7 slides
    num_slides = len(prs.slides)
    if num_slides != 7:
        print(f"CRITICAL: Expected 7 slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Loaded presentation with {num_slides} slides")

    # Collect all slide notes
    slide_notes = []
    for i, slide in enumerate(prs.slides):
        try:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            notes_text = ""
        slide_notes.append(notes_text)
        print(f"  Slide {i+1} notes preview: {repr(notes_text[:80] if notes_text else '')}")

    # Component 1: All 7 slides have non-empty notes (0.2 points)
    # This FAILS on initial (all empty) → PASSES on golden (all populated)
    try:
        non_empty_count = sum(1 for n in slide_notes if n.strip())
        if non_empty_count == 7:
            print(f"PASS: Component 1 — All 7 slides have non-empty notes (0.2 pts)")
            total_score += 0.2
        elif non_empty_count >= 5:
            print(f"PARTIAL: Component 1 — {non_empty_count}/7 slides have non-empty notes (0.1 pts)")
            total_score += 0.1
        else:
            print(f"FAIL: Component 1 — Only {non_empty_count}/7 slides have non-empty notes (expected 7)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: All 7 slides contain '---' separator in notes (0.3 points)
    # This FAILS on initial (all empty) → PASSES on golden (all have separator)
    try:
        separator_count = sum(1 for n in slide_notes if '---' in n)
        if separator_count == 7:
            print(f"PASS: Component 2 — All 7 slides contain '---' separator (0.3 pts)")
            total_score += 0.3
        elif separator_count >= 5:
            print(f"PARTIAL: Component 2 — {separator_count}/7 slides contain '---' separator (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 — Only {separator_count}/7 slides contain '---' separator")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: English text correctly placed (before '---') in >= 5/7 slides (0.25 points)
    # This FAILS on initial (no notes) → PASSES on golden (English text is before '---')
    try:
        en_correct = 0
        for i, notes in enumerate(slide_notes):
            if '---' not in notes:
                continue
            # Get text before separator
            before_sep = notes.split('---')[0].strip()
            expected_en = ENGLISH_NOTES[i].strip()
            if before_sep == expected_en:
                en_correct += 1
            else:
                # Try a relaxed check: expected English text is contained in the before_sep portion
                if expected_en in before_sep or before_sep in expected_en:
                    en_correct += 1
                    print(f"  Slide {i+1}: English text matches (relaxed check)")
                else:
                    print(f"  Slide {i+1}: English mismatch. Before '---': {repr(before_sep[:80])}")

        if en_correct == 7:
            print(f"PASS: Component 3 — English notes correct in all 7 slides (0.25 pts)")
            total_score += 0.25
        elif en_correct >= 5:
            partial = round(0.25 * en_correct / 7, 4)
            print(f"PARTIAL: Component 3 — English notes correct in {en_correct}/7 slides ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — English notes correct in only {en_correct}/7 slides")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: French text correctly placed (after '---') in >= 5/7 slides (0.25 points)
    # This FAILS on initial (no notes) → PASSES on golden (French text is after '---')
    try:
        fr_correct = 0
        for i, notes in enumerate(slide_notes):
            if '---' not in notes:
                continue
            # Get text after separator
            after_sep = notes.split('---', 1)[1].strip()
            expected_fr = FRENCH_NOTES[i].strip()
            if after_sep == expected_fr:
                fr_correct += 1
            else:
                # Try a relaxed check: expected French text is contained in the after_sep portion
                if expected_fr in after_sep or after_sep in expected_fr:
                    fr_correct += 1
                    print(f"  Slide {i+1}: French text matches (relaxed check)")
                else:
                    print(f"  Slide {i+1}: French mismatch. After '---': {repr(after_sep[:80])}")

        if fr_correct == 7:
            print(f"PASS: Component 4 — French notes correct in all 7 slides (0.25 pts)")
            total_score += 0.25
        elif fr_correct >= 5:
            partial = round(0.25 * fr_correct / 7, 4)
            print(f"PARTIAL: Component 4 — French notes correct in {fr_correct}/7 slides ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — French notes correct in only {fr_correct}/7 slides")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = round(min(total_score, 1.0), 4)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
if not os.path.exists(PPTX_PATH):
    print(f"File not found: {PPTX_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(PPTX_PATH)
