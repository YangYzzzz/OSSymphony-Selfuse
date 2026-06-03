"""
Initial Setup: French language class presentation with separate notes documents
Task ID: osworld_multi_apps_impress_notes_import_013
Domain: libreoffice_impress

Creates:
  - French_Lesson.pptx (7 slides, empty notes) on Desktop
  - french_notes_en.docx (English notes for 7 slides) on Desktop
  - french_notes_fr.docx (French notes for 7 slides) on Desktop

The task: insert both sets of notes into each slide (English first, then ---, then French).
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from docx import Document

DESKTOP = '/home/user/Desktop'
PPTX_PATH = f'{DESKTOP}/French_Lesson.pptx'
EN_NOTES_PATH = f'{DESKTOP}/french_notes_en.docx'
FR_NOTES_PATH = f'{DESKTOP}/french_notes_fr.docx'


# ---------------------------------------------------------------------------
# Slide content: (title, subtitle_or_body)
# ---------------------------------------------------------------------------
SLIDE_CONTENT = [
    {
        "title": "Introduction to French",
        "body": "Welcome to Beginner French\nLesson 1: Greetings and Introductions",
    },
    {
        "title": "Common Greetings",
        "body": "• Bonjour – Good morning / Hello\n• Bonsoir – Good evening\n• Salut – Hi (informal)\n• Au revoir – Goodbye\n• À bientôt – See you soon",
    },
    {
        "title": "Introducing Yourself",
        "body": "• Je m'appelle... – My name is...\n• J'ai ... ans – I am ... years old\n• Je suis de... – I am from...\n• Enchanté(e) – Nice to meet you",
    },
    {
        "title": "Numbers 1–10",
        "body": "Un, Deux, Trois, Quatre, Cinq\nSix, Sept, Huit, Neuf, Dix",
    },
    {
        "title": "Colors in French",
        "body": "• Rouge – Red\n• Bleu – Blue\n• Vert – Green\n• Jaune – Yellow\n• Blanc – White\n• Noir – Black",
    },
    {
        "title": "Days of the Week",
        "body": "Lundi • Mardi • Mercredi\nJeudi • Vendredi\nSamedi • Dimanche",
    },
    {
        "title": "Classroom Phrases",
        "body": "• Répétez, s'il vous plaît – Please repeat\n• Je ne comprends pas – I don't understand\n• Comment dit-on... ? – How do you say...?\n• Pouvez-vous parler plus lentement? – Can you speak more slowly?",
    },
]

# ---------------------------------------------------------------------------
# Notes content (English) — per slide
# ---------------------------------------------------------------------------
EN_NOTES = [
    # Slide 1
    "Welcome students to the first French lesson. Explain the course structure: "
    "7 sessions covering greetings, numbers, colors, days, and classroom language. "
    "Ask students about their motivation for learning French.",

    # Slide 2
    "Drill each greeting with the class. Emphasize that 'Bonjour' is used until "
    "approximately 6 PM, after which 'Bonsoir' is more appropriate. Note that 'Salut' "
    "is informal and should not be used with strangers or teachers.",

    # Slide 3
    "Have each student introduce themselves using the phrases on screen. Remind them "
    "that 'Enchanté' is masculine and 'Enchantée' is feminine. Practice: Je m'appelle "
    "[Name], j'ai [age] ans, je suis de [city].",

    # Slide 4
    "Count together as a class from 1 to 10 and back. Point out the pronunciation "
    "differences: 'cinq' is /sɛ̃k/, 'six' is /sis/ in isolation but /si/ before a "
    "consonant. Use flashcards for a quick memory game.",

    # Slide 5
    "Show color swatches alongside the French words. Note that adjectives in French "
    "agree in gender and number with the noun. Example: 'un ballon rouge' (masculine) "
    "vs 'une voiture rouge' (feminine). Colors ending in 'e' don't change form.",

    # Slide 6
    "Note that the French week starts on Monday, not Sunday. Lundi through Vendredi "
    "are weekdays; Samedi and Dimanche are the weekend. Ask students to name their "
    "favorite day in French and explain why.",

    # Slide 7
    "Practice each classroom phrase aloud. Encourage students to use these expressions "
    "throughout the rest of the course. Stress that asking for clarification is a sign "
    "of engagement, not weakness. End the lesson with a brief review quiz.",
]

# ---------------------------------------------------------------------------
# Notes content (French) — per slide
# ---------------------------------------------------------------------------
FR_NOTES = [
    # Slide 1
    "Bienvenue aux étudiants dans ce premier cours de français. Expliquez la structure "
    "du cours : 7 séances couvrant les salutations, les chiffres, les couleurs, les "
    "jours et le langage de classe. Demandez aux étudiants leur motivation pour apprendre le français.",

    # Slide 2
    "Pratiquez chaque salutation avec la classe. Soulignez que 'Bonjour' s'utilise "
    "jusqu'à environ 18h, après quoi 'Bonsoir' est plus approprié. Notez que 'Salut' "
    "est informel et ne doit pas être utilisé avec des inconnus ou des professeurs.",

    # Slide 3
    "Demandez à chaque étudiant de se présenter en utilisant les expressions à l'écran. "
    "Rappelez-leur qu''Enchanté' est masculin et 'Enchantée' est féminin. Pratique : "
    "Je m'appelle [Prénom], j'ai [âge] ans, je suis de [ville].",

    # Slide 4
    "Comptez ensemble de 1 à 10 et retour. Signalez les différences de prononciation : "
    "'cinq' se prononce /sɛ̃k/, 'six' se prononce /sis/ seul mais /si/ devant une "
    "consonne. Utilisez des cartes mémoire pour un jeu rapide de mémorisation.",

    # Slide 5
    "Montrez des échantillons de couleurs avec les mots français. Notez que les adjectifs "
    "en français s'accordent en genre et en nombre avec le nom. Exemple : 'un ballon rouge' "
    "(masculin) vs 'une voiture rouge' (féminin). Les couleurs se terminant par 'e' ne changent pas de forme.",

    # Slide 6
    "Notez que la semaine française commence le lundi, pas le dimanche. Du lundi au "
    "vendredi, ce sont les jours de semaine ; samedi et dimanche sont le week-end. "
    "Demandez aux étudiants de nommer leur jour préféré en français et d'expliquer pourquoi.",

    # Slide 7
    "Pratiquez chaque expression de classe à voix haute. Encouragez les étudiants à "
    "utiliser ces expressions tout au long du reste du cours. Insistez sur le fait que "
    "demander des éclaircissements est un signe d'engagement, pas de faiblesse. "
    "Terminez la leçon par un quiz de révision rapide.",
]


def create_pptx():
    """Create French_Lesson.pptx with 7 slides and empty notes."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for i, content in enumerate(SLIDE_CONTENT):
        if i == 0:
            layout = prs.slide_layouts[0]  # Title Slide
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = content["title"]
            slide.placeholders[1].text = content["body"]
        else:
            layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = content["title"]
            slide.placeholders[1].text = content["body"]
        # Notes must be EMPTY in initial state
        # Do NOT access slide.notes_slide to avoid creating a notes pane with content.
        # Leave notes untouched (they are empty by default).

    prs.save(PPTX_PATH)
    print(f'Created: {PPTX_PATH}')


def create_en_notes_docx():
    """Create french_notes_en.docx with English notes for each slide."""
    doc = Document()
    doc.add_heading('French Lesson – Speaker Notes (English)', level=1)

    for idx, note in enumerate(EN_NOTES, 1):
        doc.add_heading(f'Slide {idx}', level=2)
        doc.add_paragraph(note)

    doc.save(EN_NOTES_PATH)
    print(f'Created: {EN_NOTES_PATH}')


def create_fr_notes_docx():
    """Create french_notes_fr.docx with French notes for each slide."""
    doc = Document()
    doc.add_heading('Leçon de Français – Notes du Présentateur (Français)', level=1)

    for idx, note in enumerate(FR_NOTES, 1):
        doc.add_heading(f'Diapositive {idx}', level=2)
        doc.add_paragraph(note)

    doc.save(FR_NOTES_PATH)
    print(f'Created: {FR_NOTES_PATH}')


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch a GUI app on the VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def main():
    os.makedirs(DESKTOP, exist_ok=True)

    create_pptx()
    create_en_notes_docx()
    create_fr_notes_docx()

    # Open French_Lesson.pptx in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{PPTX_PATH}"', delay_sec=3.0)

    print('GUI_READY: launched LibreOffice Impress with French_Lesson.pptx (DISPLAY=:0)')


main()
