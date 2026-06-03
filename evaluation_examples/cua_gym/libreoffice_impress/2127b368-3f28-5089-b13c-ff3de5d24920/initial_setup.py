"""
Initial Setup: Add alt text to images on slide 4 of Cell Biology presentation
Task ID: impress_teach_037
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from PIL import Image, ImageDraw, ImageFont

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_037'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_placeholder_image(path, width, height, color, label):
    """Create a simple colored image with a label."""
    img = Image.new('RGB', (width, height), color)
    draw = ImageDraw.Draw(img)
    # Draw label text centered
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 20)
    except (IOError, OSError):
        font = ImageFont.load_default()
    bbox = draw.textbbox((0, 0), label, font=font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    x = (width - tw) // 2
    y = (height - th) // 2
    draw.text((x, y), label, fill=(255, 255, 255), font=font)
    img.save(path)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Cell Biology"
    slide1.placeholders[1].text = "An Introduction to the Building Blocks of Life"

    # --- Slide 2: What is a Cell? ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "What is a Cell?"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Cells are the fundamental units of life"
    p = body2.add_paragraph()
    p.text = "All living organisms are composed of one or more cells"
    p.level = 1
    p = body2.add_paragraph()
    p.text = "Cells arise from pre-existing cells (cell theory)"
    p.level = 1
    p = body2.add_paragraph()
    p.text = "The average human body contains approximately 37.2 trillion cells"
    p.level = 1

    # --- Slide 3: Types of Cells ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Types of Cells"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Prokaryotic Cells"
    p = body3.add_paragraph()
    p.text = "No membrane-bound nucleus; bacteria and archaea"
    p.level = 1
    p = body3.add_paragraph()
    p.text = "Eukaryotic Cells"
    p = body3.add_paragraph()
    p.text = "Contain a nucleus and organelles; plants, animals, fungi"
    p.level = 1

    # --- Slide 4: Plant Cell Structure (3 images, NO alt text) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Title textbox
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Plant Cell Structure & Photosynthesis"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Create 3 placeholder images
    img1_path = f'{WORKDIR}/img_plant_cell.png'
    img2_path = f'{WORKDIR}/img_leaf_cross_section.png'
    img3_path = f'{WORKDIR}/img_photosynthesis_chart.png'

    create_placeholder_image(img1_path, 400, 300, (46, 125, 50), "Plant Cell Diagram")
    create_placeholder_image(img2_path, 400, 300, (100, 65, 120), "Leaf Cross-Section")
    create_placeholder_image(img3_path, 400, 300, (30, 80, 140), "Photosynthesis Rate Chart")

    # Add images to slide 4 - explicitly clear alt text (python-pptx defaults descr to filename)
    pic1 = slide4.shapes.add_picture(img1_path, Inches(0.5), Inches(1.3), Inches(3.8), Inches(2.8))
    pic2 = slide4.shapes.add_picture(img2_path, Inches(4.7), Inches(1.3), Inches(3.8), Inches(2.8))
    pic3 = slide4.shapes.add_picture(img3_path, Inches(9.0), Inches(1.3), Inches(3.8), Inches(2.8))

    # Clear alt text on all three images
    for pic in [pic1, pic2, pic3]:
        cNvPr = pic._element.find(qn('p:nvPicPr') + '/' + qn('p:cNvPr'))
        if cNvPr is not None:
            if 'descr' in cNvPr.attrib:
                del cNvPr.attrib['descr']

    # Add captions below images
    for x, caption in [(0.5, "Figure 1: Plant Cell"), (4.7, "Figure 2: Leaf Cross-Section"), (9.0, "Figure 3: Photosynthesis Rate")]:
        cap = slide4.shapes.add_textbox(Inches(x), Inches(4.3), Inches(3.8), Inches(0.5))
        cp = cap.text_frame.paragraphs[0]
        cp.text = caption
        cp.alignment = PP_ALIGN.CENTER
        cp.runs[0].font.size = Pt(11)
        cp.runs[0].font.italic = True

    # --- Slide 5: Organelles Overview ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Key Organelles"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Nucleus – Contains genetic material (DNA)"
    for item in [
        "Mitochondria – Powerhouse of the cell (ATP production)",
        "Chloroplasts – Site of photosynthesis (plant cells)",
        "Endoplasmic Reticulum – Protein and lipid synthesis",
        "Golgi Apparatus – Packaging and transport of proteins",
        "Ribosomes – Protein synthesis from mRNA",
    ]:
        p = body5.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 6: Cell Division ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Cell Division"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Mitosis – Produces identical daughter cells for growth and repair"
    p = body6.add_paragraph()
    p.text = "Meiosis – Produces gametes with half the chromosome number"
    p = body6.add_paragraph()
    p.text = "Interphase: G1 → S → G2 (cell prepares for division)"
    p.level = 1
    p = body6.add_paragraph()
    p.text = "Mitotic Phase: Prophase → Metaphase → Anaphase → Telophase"
    p.level = 1

    # --- Slide 7: Cellular Respiration ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Cellular Respiration"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "C₆H₁₂O₆ + 6O₂ → 6CO₂ + 6H₂O + ATP"
    p = body7.add_paragraph()
    p.text = "Glycolysis (cytoplasm) – 2 ATP"
    p.level = 1
    p = body7.add_paragraph()
    p.text = "Krebs Cycle (mitochondrial matrix) – 2 ATP"
    p.level = 1
    p = body7.add_paragraph()
    p.text = "Electron Transport Chain (inner membrane) – 34 ATP"
    p.level = 1

    # --- Slide 8: Summary ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Summary & Key Takeaways"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Cells are the basic structural and functional unit of all organisms"
    p = body8.add_paragraph()
    p.text = "Plant cells have unique structures: cell wall, chloroplasts, large vacuole"
    p = body8.add_paragraph()
    p.text = "Photosynthesis converts light energy to chemical energy"
    p = body8.add_paragraph()
    p.text = "Cellular respiration releases stored energy as ATP"
    p = body8.add_paragraph()
    p.text = "Cell division ensures growth, repair, and reproduction"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Clean up temp images
    for p in [img1_path, img2_path, img3_path]:
        if os.path.exists(p):
            os.remove(p)

    # Launch LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
