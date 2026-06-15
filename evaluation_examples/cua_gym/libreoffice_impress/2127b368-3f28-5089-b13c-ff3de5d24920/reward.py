"""
Reward Script: Add alt text to three images on slide 4
Task ID: impress_teach_037
Domain: libreoffice_impress
Scoring:
  Component 1: Image 1 alt text matches expected (0.33 pts)
  Component 2: Image 2 alt text matches expected (0.33 pts)
  Component 3: Image 3 alt text matches expected (0.34 pts)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_037'

# Expected alt text for each image on slide 4 (in order of appearance)
EXPECTED_ALT_TEXTS = [
    "Diagram of plant cell showing chloroplasts and nucleus",
    "Microscope photograph of leaf cross-section",
    "Chart showing photosynthesis rate vs light intensity",
]

COMPONENT_WEIGHTS = [0.33, 0.33, 0.34]


def persist_app_state(domain: str):
    """Attempt to save any unsaved GUI state via Ctrl+S."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_picture_alt_text(shape):
    """Extract the alt text (descr attribute) from a picture shape's cNvPr element."""
    ns = {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }
    cNvPr = shape._element.find('.//p:nvPicPr/p:cNvPr', ns)
    if cNvPr is not None:
        return cNvPr.get('descr', '')
    return ''


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, expected at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # slide 4 (0-indexed)

    # Collect all picture shapes on slide 4
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]

    # Precondition: exactly 3 pictures on slide 4
    if len(pictures) < 3:
        print(f"FAIL: Slide 4 has {len(pictures)} pictures, expected 3")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Image 1 alt text (0.33 points)
    try:
        alt1 = get_picture_alt_text(pictures[0]).strip()
        expected1 = EXPECTED_ALT_TEXTS[0]
        if alt1.lower() == expected1.lower():
            print(f"PASS: Component 1 — Image 1 alt text matches: '{alt1}' (0.33 pts)")
            total_score += COMPONENT_WEIGHTS[0]
        else:
            print(f"FAIL: Component 1 — Image 1 alt text: expected '{expected1}', found '{alt1}'")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Image 2 alt text (0.33 points)
    try:
        alt2 = get_picture_alt_text(pictures[1]).strip()
        expected2 = EXPECTED_ALT_TEXTS[1]
        if alt2.lower() == expected2.lower():
            print(f"PASS: Component 2 — Image 2 alt text matches: '{alt2}' (0.33 pts)")
            total_score += COMPONENT_WEIGHTS[1]
        else:
            print(f"FAIL: Component 2 — Image 2 alt text: expected '{expected2}', found '{alt2}'")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Image 3 alt text (0.34 points)
    try:
        alt3 = get_picture_alt_text(pictures[2]).strip()
        expected3 = EXPECTED_ALT_TEXTS[2]
        if alt3.lower() == expected3.lower():
            print(f"PASS: Component 3 — Image 3 alt text matches: '{alt3}' (0.34 pts)")
            total_score += COMPONENT_WEIGHTS[2]
        else:
            print(f"FAIL: Component 3 — Image 3 alt text: expected '{expected3}', found '{alt3}'")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist any unsaved GUI state before verification
persist_app_state("libreoffice_impress")

# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
