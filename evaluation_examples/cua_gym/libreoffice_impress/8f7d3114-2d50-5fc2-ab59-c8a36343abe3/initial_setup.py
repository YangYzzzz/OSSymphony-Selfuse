"""
Initial Setup: Create a 10-slide presentation with white bg on slide 7,
and a 100x100 pattern image at ~/Pictures/small_pattern.png.
Task ID: impress_el_074
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

WORKDIR = '/home/user'
TASK_ID = 'impress_el_074'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
PATTERN_DIR = f'{WORKDIR}/Pictures'
PATTERN_PATH = f'{PATTERN_DIR}/small_pattern.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_pattern_image():
    """Create a 100x100 px checkerboard-style pattern image."""
    os.makedirs(PATTERN_DIR, exist_ok=True)
    img = Image.new('RGB', (100, 100), (255, 255, 255))
    pixels = img.load()
    # Create a subtle geometric pattern: diagonal stripes
    for y in range(100):
        for x in range(100):
            if (x + y) % 20 < 10:
                pixels[x, y] = (70, 130, 180)  # steel blue
            else:
                pixels[x, y] = (245, 245, 245)  # near white
    img.save(PATTERN_PATH)
    print(f'Pattern image created: {PATTERN_PATH}')


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_contents = [
        {
            'title': 'Q3 2025 Strategic Review',
            'subtitle': 'Meridian Consulting Group',
            'layout': 0,
        },
        {
            'title': 'Executive Summary',
            'body': (
                'Revenue grew 18% year-over-year driven by expansion in APAC markets.\n'
                'Client retention rate improved to 94.2%, surpassing the 92% target.\n'
                'Three new enterprise accounts were onboarded in the healthcare vertical.\n'
                'Operating margin held steady at 23.5% despite increased hiring.'
            ),
            'layout': 1,
        },
        {
            'title': 'Market Landscape',
            'body': (
                'The global consulting market reached $320B in 2025.\n'
                'Digital transformation services accounted for 42% of new engagements.\n'
                'AI-driven analytics became a key differentiator for top-tier firms.\n'
                'Regulatory compliance work surged 25% in financial services.'
            ),
            'layout': 1,
        },
        {
            'title': 'Client Portfolio Breakdown',
            'body': (
                'Financial Services: 34% of revenue ($12.8M)\n'
                'Healthcare & Life Sciences: 22% of revenue ($8.3M)\n'
                'Technology: 19% of revenue ($7.1M)\n'
                'Energy & Utilities: 15% of revenue ($5.6M)\n'
                'Public Sector: 10% of revenue ($3.7M)'
            ),
            'layout': 1,
        },
        {
            'title': 'Team Growth & Development',
            'body': (
                'Headcount increased from 185 to 214 employees.\n'
                'Promoted 12 senior consultants to principal level.\n'
                'Launched the Emerging Leaders Program with 30 participants.\n'
                'Average training hours per employee: 48 hours in Q3.'
            ),
            'layout': 1,
        },
        {
            'title': 'Technology Investments',
            'body': (
                'Deployed proprietary analytics platform "Meridian Insight" to 85% of projects.\n'
                'Migrated internal infrastructure to hybrid cloud architecture.\n'
                'Integrated GenAI assistants into proposal generation workflow.\n'
                'Reduced project setup time by 35% through automation.'
            ),
            'layout': 1,
        },
        {
            # Slide 7 (index 6) — white/default background, no tiled image
            'title': 'Regional Performance Overview',
            'body': (
                'North America: $21.4M revenue, +12% YoY\n'
                'Europe: $9.8M revenue, +8% YoY\n'
                'Asia-Pacific: $5.2M revenue, +31% YoY\n'
                'Latin America: $1.1M revenue, +15% YoY'
            ),
            'layout': 1,
        },
        {
            'title': 'Key Risks & Mitigation',
            'body': (
                'Talent attrition in data science roles — expanding remote work options.\n'
                'Currency fluctuations impacting EMEA margins — hedging strategy in place.\n'
                'Regulatory changes in EU AI Act — compliance task force established.\n'
                'Client budget freezes in tech sector — diversifying pipeline focus.'
            ),
            'layout': 1,
        },
        {
            'title': 'Q4 2025 Priorities',
            'body': (
                '1. Close $8M pipeline in healthcare vertical by December.\n'
                '2. Launch Meridian Insight v2.0 with predictive modeling.\n'
                '3. Expand Singapore office to 25 consultants.\n'
                '4. Achieve ISO 27001 certification for data handling practices.'
            ),
            'layout': 1,
        },
        {
            'title': 'Thank You',
            'subtitle': 'Questions & Discussion\ncontact@meridianconsulting.com',
            'layout': 0,
        },
    ]

    for i, sc in enumerate(slide_contents):
        layout_idx = sc['layout']
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        if slide.shapes.title:
            slide.shapes.title.text = sc['title']
            for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                run.font.size = Pt(32)
                run.font.bold = True

        if 'subtitle' in sc and len(slide.placeholders) > 1:
            slide.placeholders[1].text = sc['subtitle']
        elif 'body' in sc and len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            lines = sc['body'].split('\n')
            for j, line in enumerate(lines):
                if j == 0:
                    tf.paragraphs[0].text = line
                    tf.paragraphs[0].space_after = Pt(6)
                else:
                    p = tf.add_paragraph()
                    p.text = line
                    p.space_after = Pt(6)
                para = tf.paragraphs[j]
                for run in para.runs:
                    run.font.size = Pt(18)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


# Create pattern image first, then presentation
create_pattern_image()
create_initial()
