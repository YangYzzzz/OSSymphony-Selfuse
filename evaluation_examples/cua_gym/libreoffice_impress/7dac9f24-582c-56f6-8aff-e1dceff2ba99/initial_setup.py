"""
Initial Setup: Sales Funnel Presentation with empty slide 5
Task ID: impress_sales_047
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_047'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_placeholder(shape, text, font_size=Pt(18), bold=False, color=None):
    """Helper to set text on a placeholder or text frame."""
    tf = shape.text_frame
    tf.paragraphs[0].text = text
    for run in tf.paragraphs[0].runs:
        run.font.size = font_size
        run.font.bold = bold
        if color:
            run.font.color.rgb = color


def add_bullet_slide(prs, layout, title_text, bullets):
    """Add a slide with title and bullet points."""
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.level = 0
    return slide


def create_initial():
    prs = Presentation()
    # Standard widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]      # Title Slide
    content_layout = prs.slide_layouts[1]     # Title + Content
    blank_layout = prs.slide_layouts[5]       # Blank
    title_only_layout = prs.slide_layouts[5]  # We'll use blank and add title manually

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(title_layout)
    slide1.shapes.title.text = "Q4 Sales Strategy"
    slide1.placeholders[1].text = "Accelerating Growth Through Targeted Engagement"

    # --- Slide 2: Market Overview ---
    slide2 = add_bullet_slide(prs, content_layout, "Market Overview", [
        "Total addressable market grew 12% YoY to $4.8B",
        "Enterprise segment accounts for 58% of revenue",
        "Mid-market adoption accelerating in APAC region",
        "Competitive landscape shifting toward SaaS-first solutions",
        "Customer acquisition cost decreased 8% through digital channels",
    ])

    # --- Slide 3: Revenue Targets ---
    slide3 = prs.slides.add_slide(content_layout)
    slide3.shapes.title.text = "Revenue Targets"
    # Add a table
    rows, cols = 5, 4
    tbl_shape = slide3.shapes.add_table(
        rows, cols, Inches(1.5), Inches(2), Inches(10), Inches(4)
    )
    table = tbl_shape.table
    headers = ["Region", "Q3 Actual ($M)", "Q4 Target ($M)", "Growth %"]
    data = [
        ["North America", "$12.4", "$14.8", "19.4%"],
        ["Europe", "$8.7", "$10.2", "17.2%"],
        ["Asia-Pacific", "$5.3", "$6.9", "30.2%"],
        ["Latin America", "$2.1", "$2.8", "33.3%"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 4: Team Structure ---
    slide4 = add_bullet_slide(prs, content_layout, "Team Structure", [
        "VP of Sales: Katherine Reyes (15 years experience)",
        "Enterprise Team: 24 account executives across 3 regions",
        "Mid-Market Team: 18 reps focused on $50K-$500K deals",
        "Sales Engineering: 12 specialists for technical demos",
        "Channel Partners: 45 active reseller relationships",
    ])

    # --- Slide 5: Sales Funnel (EMPTY content area - title only) ---
    slide5 = prs.slides.add_slide(blank_layout)
    # Add title manually as a text box
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Sales Funnel"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    # Content area intentionally left empty - agent must add pyramid diagram

    # --- Slide 6: Action Items ---
    slide6 = add_bullet_slide(prs, content_layout, "Action Items", [
        "Launch targeted email campaign for enterprise prospects by Oct 15",
        "Schedule quarterly business reviews with top 20 accounts",
        "Deploy updated CRM dashboards for pipeline visibility",
        "Conduct competitive analysis refresh for Q4 positioning",
        "Finalize partner incentive program for holiday season push",
    ])

    # --- Slide 7: Timeline ---
    slide7 = add_bullet_slide(prs, content_layout, "Implementation Timeline", [
        "October 1-15: Campaign preparation and asset creation",
        "October 16-31: Enterprise outreach and demo scheduling",
        "November 1-15: Mid-market blitz and partner activation",
        "November 16-30: Pipeline review and forecast adjustments",
        "December 1-31: Close push and year-end deal acceleration",
    ])

    # --- Slide 8: Thank You ---
    slide8 = prs.slides.add_slide(title_layout)
    slide8.shapes.title.text = "Thank You"
    slide8.placeholders[1].text = "Questions? Contact sales-strategy@acmecorp.com"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
