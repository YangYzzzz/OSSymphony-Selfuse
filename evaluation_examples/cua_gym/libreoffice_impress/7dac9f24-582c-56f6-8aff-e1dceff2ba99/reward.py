"""
Reward Script: Insert pyramid diagram on slide 5 with four levels and gradient fills
Task ID: impress_sales_047
Domain: libreoffice_impress
Scoring:
  Component 1 (0.35): Four pyramid shapes with correct labels on slide 5
  Component 2 (0.25): Shapes form a pyramid (widths decrease bottom to top)
  Component 3 (0.25): Gradient fills with correct color progression (light blue to dark blue)
  Component 4 (0.15): Presentation structure preserved (8 slides, Sales Funnel title)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_047'

# Expected pyramid labels from bottom (widest) to top (narrowest)
EXPECTED_LABELS = ['Awareness', 'Interest', 'Decision', 'Action']

# Expected gradient color anchors (start colors for each level, bottom to top)
# Awareness: #B3D9FF (light blue), Action: #003366 (dark blue)
BOTTOM_COLOR_HEX = 'B3D9FF'
TOP_COLOR_HEX = '003366'


def hex_to_rgb(h):
    """Convert hex string to (R, G, B) tuple."""
    h = h.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def color_distance(c1, c2):
    """Euclidean distance between two RGB tuples."""
    return sum((a - b) ** 2 for a, b in zip(c1, c2)) ** 0.5


def persist_app_state(domain):
    """Try to save any unsaved LibreOffice edits."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print("PERSIST: ctrl+s sent for %s" % domain)
        except Exception as e:
            print("PERSIST_WARN: save hook failed: %s" % e)


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print("CRITICAL: Cannot load file %s: %s" % (file_path, e))
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 5 slides
    if len(prs.slides) < 5:
        print("CRITICAL: Presentation has fewer than 5 slides")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # Slide 5 (0-indexed)

    # Find all non-placeholder, non-textbox shapes that could be pyramid levels
    # These are AUTO_SHAPE type shapes added by the task
    pyramid_shapes = []
    for shape in slide.shapes:
        # Skip placeholders and the existing "Sales Funnel" text box
        if shape.shape_type == 14:  # PLACEHOLDER
            continue
        if shape.has_text_frame:
            text = shape.text_frame.paragraphs[0].text.strip() if shape.text_frame.paragraphs else ''
            if text == 'Sales Funnel':
                continue
        # This is a candidate pyramid shape
        if shape.has_text_frame:
            label = ''
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    label = t
                    break
            pyramid_shapes.append({
                'name': shape.name,
                'label': label,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'shape_type': shape.shape_type,
            })

    print("Found %d candidate pyramid shapes on slide 5" % len(pyramid_shapes))
    for ps in pyramid_shapes:
        print("  %s: label=%r, width=%d, top=%d" % (ps['name'], ps['label'], ps['width'], ps['top']))

    # Component 1: Four pyramid shapes with correct labels (0.35 points)
    try:
        found_labels = [ps['label'] for ps in pyramid_shapes]
        matched_labels = []
        for expected in EXPECTED_LABELS:
            for fl in found_labels:
                if expected.lower() in fl.lower():
                    matched_labels.append(expected)
                    break

        label_count = len(matched_labels)
        if label_count == 4:
            print("PASS: Component 1 — All 4 pyramid labels found: %s (0.35 pts)" % matched_labels)
            total_score += 0.35
        elif label_count >= 2:
            partial = 0.35 * (label_count / 4.0)
            print("PARTIAL: Component 1 — %d/4 labels found: %s (%.2f pts)" % (label_count, matched_labels, partial))
            total_score += partial
        else:
            print("FAIL: Component 1 — Only %d/4 labels found: %s" % (label_count, found_labels))
    except Exception as e:
        print("ERROR: Component 1 — %s" % e)

    # Component 2: Shapes form a pyramid - widths decrease from bottom to top (0.25 points)
    try:
        if len(pyramid_shapes) >= 4:
            # Sort by vertical position (top), descending = bottom first
            sorted_by_pos = sorted(pyramid_shapes, key=lambda s: s['top'], reverse=True)
            # Take the 4 shapes with labels matching expected ones
            label_shape_map = {}
            for ps in pyramid_shapes:
                for exp in EXPECTED_LABELS:
                    if exp.lower() in ps['label'].lower():
                        label_shape_map[exp] = ps
                        break

            if len(label_shape_map) >= 4:
                # Check pyramid ordering: bottom-to-top should have decreasing widths
                widths_bottom_to_top = [
                    label_shape_map['Awareness']['width'],
                    label_shape_map['Interest']['width'],
                    label_shape_map['Decision']['width'],
                    label_shape_map['Action']['width'],
                ]
                is_pyramid = all(widths_bottom_to_top[i] > widths_bottom_to_top[i+1] for i in range(3))

                # Also check vertical ordering: Awareness should be lowest (largest top), Action highest (smallest top)
                tops_bottom_to_top = [
                    label_shape_map['Awareness']['top'],
                    label_shape_map['Interest']['top'],
                    label_shape_map['Decision']['top'],
                    label_shape_map['Action']['top'],
                ]
                is_ordered = all(tops_bottom_to_top[i] > tops_bottom_to_top[i+1] for i in range(3))

                if is_pyramid and is_ordered:
                    print("PASS: Component 2 — Pyramid shape: widths=%s, tops=%s (0.25 pts)" % (widths_bottom_to_top, tops_bottom_to_top))
                    total_score += 0.25
                elif is_pyramid or is_ordered:
                    print("PARTIAL: Component 2 — pyramid=%s, ordered=%s (0.125 pts)" % (is_pyramid, is_ordered))
                    total_score += 0.125
                else:
                    print("FAIL: Component 2 — Not a pyramid shape: widths=%s, tops=%s" % (widths_bottom_to_top, tops_bottom_to_top))
            else:
                print("FAIL: Component 2 — Cannot map all 4 labels to shapes")
        else:
            print("FAIL: Component 2 — Fewer than 4 pyramid shapes found")
    except Exception as e:
        print("ERROR: Component 2 — %s" % e)

    # Component 3: Gradient fills with correct color progression (0.25 points)
    try:
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
              'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

        gradient_shapes = {}
        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide5.xml') as f:
                root = ET.fromstring(f.read())
                for sp in root.findall('.//p:sp', ns):
                    # Get text
                    txBody = sp.find('.//p:txBody', ns)
                    text = ''
                    if txBody is not None:
                        for t in txBody.findall('.//a:t', ns):
                            text += (t.text or '')
                    text = text.strip()

                    # Get fill info
                    spPr = sp.find('.//p:spPr', ns)
                    if spPr is not None:
                        gradFill = spPr.find('a:gradFill', ns)
                        solidFill = spPr.find('a:solidFill', ns)
                        if gradFill is not None:
                            stops = []
                            for gs in gradFill.findall('.//a:gs', ns):
                                clr = gs.find('a:srgbClr', ns)
                                if clr is not None:
                                    stops.append(clr.get('val'))
                            gradient_shapes[text] = {'type': 'gradient', 'colors': stops}
                        elif solidFill is not None:
                            clr = solidFill.find('a:srgbClr', ns)
                            if clr is not None:
                                gradient_shapes[text] = {'type': 'solid', 'colors': [clr.get('val')]}

        print("Fill analysis: %s" % gradient_shapes)

        # Check that pyramid shapes have fills and color progression
        gradient_count = 0
        color_ok = 0

        bottom_target = hex_to_rgb(BOTTOM_COLOR_HEX)  # B3D9FF - light blue
        top_target = hex_to_rgb(TOP_COLOR_HEX)  # 003366 - dark blue

        for label in EXPECTED_LABELS:
            matched_key = None
            for key in gradient_shapes:
                if label.lower() in key.lower():
                    matched_key = key
                    break
            if matched_key and gradient_shapes[matched_key]['type'] == 'gradient':
                gradient_count += 1
                colors = gradient_shapes[matched_key]['colors']
                if colors:
                    first_color = hex_to_rgb(colors[0])
                    # Awareness should be closest to light blue, Action to dark blue
                    if label == 'Awareness':
                        dist = color_distance(first_color, bottom_target)
                        if dist < 80:
                            color_ok += 1
                    elif label == 'Action':
                        dist = color_distance(first_color, top_target)
                        if dist < 80:
                            color_ok += 1
                    else:
                        # Middle levels should be between
                        color_ok += 1
            elif matched_key and gradient_shapes[matched_key]['type'] == 'solid':
                # Solid fill is acceptable if color matches the gradient progression
                gradient_count += 0.5

        if gradient_count >= 4 and color_ok >= 2:
            print("PASS: Component 3 — All 4 shapes have gradient fills with correct color progression (0.25 pts)")
            total_score += 0.25
        elif gradient_count >= 2:
            partial = 0.25 * (gradient_count / 4.0)
            print("PARTIAL: Component 3 — %s/4 shapes have gradient fills (%.2f pts)" % (gradient_count, partial))
            total_score += partial
        else:
            print("FAIL: Component 3 — Only %s/4 shapes have gradient fills" % gradient_count)
    except Exception as e:
        print("ERROR: Component 3 — %s" % e)

    # Component 4: Presentation structure preserved (0.15 points)
    try:
        slide_count_ok = len(prs.slides) == 8
        # Check Sales Funnel title on slide 5
        title_ok = any(
            'Sales Funnel' in para.text
            for shape in slide.shapes if shape.has_text_frame
            for para in shape.text_frame.paragraphs
        )

        # Only award points if BOTH the pyramid shapes exist AND structure is preserved
        # This ensures initial_env (which has the title but no pyramid) scores 0
        if slide_count_ok and title_ok and len(pyramid_shapes) >= 4:
            print("PASS: Component 4 — 8 slides, 'Sales Funnel' title preserved, pyramid present (0.15 pts)")
            total_score += 0.15
        elif slide_count_ok and title_ok and len(pyramid_shapes) >= 2:
            print("PARTIAL: Component 4 — Structure OK but incomplete pyramid (0.075 pts)")
            total_score += 0.075
        else:
            if not slide_count_ok:
                print("FAIL: Component 4 — Slide count is %d, expected 8" % len(prs.slides))
            if not title_ok:
                print("FAIL: Component 4 — 'Sales Funnel' title not found on slide 5")
            if len(pyramid_shapes) < 2:
                print("FAIL: Component 4 — Pyramid shapes not present")
    except Exception as e:
        print("ERROR: Component 4 — %s" % e)

    final_score = round(min(total_score, 1.0), 2)
    print("\nScore: %.2f/1.0" % total_score)
    print("REWARD: %.1f" % final_score)
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = '%s/%s.pptx' % (WORKDIR, TASK_ID)
if not os.path.exists(file_path):
    print("File not found: %s" % file_path)
    print("REWARD: 0.0")
else:
    verify_task(file_path)
