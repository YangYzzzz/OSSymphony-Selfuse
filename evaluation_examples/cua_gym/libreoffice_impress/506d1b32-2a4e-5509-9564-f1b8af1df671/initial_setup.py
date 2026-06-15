"""
Initial Setup: Fix broken hyperlink in Research Findings presentation
Task ID: impress_fix_048
Domain: libreoffice_impress

Creates a 15-slide Research Findings presentation. On slide 3, the text
'See Appendix' has an external hyperlink to http://example.com (broken),
which should instead be an internal link to slide 10.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_048'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_external_hyperlink(run, url):
    """Add an external hyperlink to a run."""
    rPr = run._r.get_or_add_rPr()
    # Create the relationship
    part = run.part
    rId = part.relate_to(url, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink', is_external=True)
    hlinkClick = rPr.makeelement(qn('a:hlinkClick'), {})
    hlinkClick.set(qn('r:id'), rId)
    rPr.append(hlinkClick)


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                alignment=PP_ALIGN.LEFT, color=None):
    """Helper to add a textbox with styled text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    for run in p.runs:
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only
    add_textbox(slide1, 1.5, 1.5, 10, 1.5, "Research Findings Report",
                font_size=36, bold=True, alignment=PP_ALIGN.CENTER,
                color=RGBColor(0x1A, 0x3C, 0x6D))
    add_textbox(slide1, 1.5, 3.5, 10, 1, "Quarterly Analysis — Q4 2025",
                font_size=20, alignment=PP_ALIGN.CENTER,
                color=RGBColor(0x4A, 0x4A, 0x4A))
    add_textbox(slide1, 1.5, 5.0, 10, 0.6, "Prepared by: Dr. Elena Vasquez, Research Division",
                font_size=14, alignment=PP_ALIGN.CENTER,
                color=RGBColor(0x70, 0x70, 0x70))

    # --- Slide 2: Table of Contents ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide2, 0.8, 0.5, 10, 1, "Table of Contents",
                font_size=28, bold=True, color=RGBColor(0x1A, 0x3C, 0x6D))
    toc_items = [
        "1. Executive Summary",
        "2. Methodology Overview",
        "3. Key Findings",
        "4. Market Analysis",
        "5. Consumer Behavior Trends",
        "6. Regional Breakdown",
        "7. Competitive Landscape",
        "8. Financial Projections",
        "9. Recommendations",
        "10. Appendix — Supplementary Data",
    ]
    toc_text = "\n".join(toc_items)
    add_textbox(slide2, 1.2, 1.8, 10, 5, toc_text, font_size=16,
                color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 3: Key Findings (with broken hyperlink) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide3, 0.8, 0.5, 10, 1, "Key Findings",
                font_size=28, bold=True, color=RGBColor(0x1A, 0x3C, 0x6D))

    findings_text = (
        "Our research indicates a significant shift in consumer preferences "
        "across all demographic segments. The primary drivers include:\n\n"
        "• Digital adoption rates increased by 34% year-over-year\n"
        "• Customer satisfaction scores improved from 72.3 to 86.1\n"
        "• Average transaction value grew by $12.40 per customer\n"
        "• Mobile engagement surpassed desktop for the first time"
    )
    add_textbox(slide3, 0.8, 1.6, 10, 3.5, findings_text, font_size=14,
                color=RGBColor(0x33, 0x33, 0x33))

    # Add the "See Appendix" text with broken external hyperlink
    txBox = slide3.shapes.add_textbox(Inches(0.8), Inches(5.5),
                                       Inches(5), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "See Appendix"
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x00, 0x56, 0xB3)
    run.font.underline = True
    # Add broken external hyperlink (should be internal to slide 10)
    add_external_hyperlink(run, "http://example.com")
    # Append explanatory text
    run2 = p.add_run()
    run2.text = " for detailed statistical tables and raw data."
    run2.font.size = Pt(14)
    run2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 4: Methodology ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide4, 0.8, 0.5, 10, 1, "Methodology",
                font_size=28, bold=True, color=RGBColor(0x1A, 0x3C, 0x6D))
    add_textbox(slide4, 0.8, 1.8, 10, 4, (
        "Data Collection Period: September 2025 — November 2025\n\n"
        "Sample Size: 4,832 respondents across 12 metropolitan areas\n\n"
        "Methods: Online surveys (62%), phone interviews (23%), "
        "focus groups (15%)\n\n"
        "Confidence Level: 95% with margin of error ±2.1%\n\n"
        "Analysis Tools: SPSS v28, R Studio, Tableau for visualization"
    ), font_size=14, color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 5: Market Analysis ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide5, 0.8, 0.5, 10, 1, "Market Analysis",
                font_size=28, bold=True, color=RGBColor(0x1A, 0x3C, 0x6D))
    add_textbox(slide5, 0.8, 1.8, 10, 4, (
        "Total addressable market grew from $4.2B to $5.8B (+38%)\n\n"
        "Key segments showing growth:\n"
        "• Enterprise solutions: +42% ($2.1B)\n"
        "• SMB platforms: +31% ($1.4B)\n"
        "• Consumer applications: +27% ($1.2B)\n"
        "• Government contracts: +18% ($0.6B)\n"
        "• Healthcare vertical: +55% ($0.5B)"
    ), font_size=14, color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 6: Consumer Behavior Trends ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide6, 0.8, 0.5, 10, 1, "Consumer Behavior Trends",
                font_size=28, bold=True, color=RGBColor(0x1A, 0x3C, 0x6D))
    add_textbox(slide6, 0.8, 1.8, 10, 4, (
        "Emerging patterns in 2025 consumer behavior:\n\n"
        "1. Preference for subscription models over one-time purchases (+28%)\n"
        "2. Increasing demand for personalized experiences (89% of respondents)\n"
        "3. Sustainability as a purchase driver (67% consider eco-impact)\n"
        "4. Social proof influence on buying decisions (74% check reviews)\n"
        "5. Voice-activated commerce adoption (31% regular users)"
    ), font_size=14, color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 7: Regional Breakdown ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide7, 0.8, 0.5, 10, 1, "Regional Breakdown",
                font_size=28, bold=True, color=RGBColor(0x1A, 0x3C, 0x6D))
    add_textbox(slide7, 0.8, 1.8, 10, 4, (
        "Performance by region (YoY growth):\n\n"
        "North America: +36% revenue, 1,842 respondents\n"
        "Europe: +29% revenue, 1,205 respondents\n"
        "Asia-Pacific: +48% revenue, 987 respondents\n"
        "Latin America: +22% revenue, 498 respondents\n"
        "Middle East & Africa: +19% revenue, 300 respondents\n\n"
        "Asia-Pacific showed the strongest momentum, driven by "
        "mobile-first adoption patterns in India and Southeast Asia."
    ), font_size=14, color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 8: Competitive Landscape ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide8, 0.8, 0.5, 10, 1, "Competitive Landscape",
                font_size=28, bold=True, color=RGBColor(0x1A, 0x3C, 0x6D))
    add_textbox(slide8, 0.8, 1.8, 10, 4, (
        "Market Share Distribution:\n\n"
        "• Vertex Technologies: 28.3% (up from 24.1%)\n"
        "• NovaCorp Solutions: 22.7% (stable)\n"
        "• Meridian Systems: 18.9% (down from 20.4%)\n"
        "• Our Company: 15.6% (up from 12.8%)\n"
        "• Others: 14.5%\n\n"
        "Notable: Our market share gain of 2.8 percentage points was "
        "the largest single-year increase among major players."
    ), font_size=14, color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 9: Financial Projections ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide9, 0.8, 0.5, 10, 1, "Financial Projections",
                font_size=28, bold=True, color=RGBColor(0x1A, 0x3C, 0x6D))
    add_textbox(slide9, 0.8, 1.8, 10, 4, (
        "Projected Revenue (in millions):\n\n"
        "FY2026 Q1: $148.2M (conservative) / $162.5M (optimistic)\n"
        "FY2026 Q2: $155.8M (conservative) / $174.3M (optimistic)\n"
        "FY2026 Q3: $163.1M (conservative) / $185.0M (optimistic)\n"
        "FY2026 Q4: $172.4M (conservative) / $198.7M (optimistic)\n\n"
        "Total FY2026: $639.5M — $720.5M\n"
        "EBITDA Margin Target: 24-27%"
    ), font_size=14, color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 10: Appendix - Supplementary Data ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide10, 0.8, 0.5, 10, 1, "Appendix — Supplementary Data",
                font_size=28, bold=True, color=RGBColor(0x1A, 0x3C, 0x6D))
    add_textbox(slide10, 0.8, 1.8, 10, 4.5, (
        "Table A1: Complete Survey Response Data\n\n"
        "Region         | Respondents | Avg. Score | Completion Rate\n"
        "North America  |   1,842     |    84.2    |     94.1%\n"
        "Europe         |   1,205     |    81.7    |     91.8%\n"
        "Asia-Pacific   |     987     |    86.9    |     96.3%\n"
        "Latin America  |     498     |    79.4    |     88.7%\n"
        "ME & Africa    |     300     |    77.1    |     85.2%\n\n"
        "Total          |   4,832     |    83.1    |     92.4%"
    ), font_size=12, color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 11: Appendix - Demographic Breakdown ---
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide11, 0.8, 0.5, 10, 1, "Appendix — Demographic Breakdown",
                font_size=28, bold=True, color=RGBColor(0x1A, 0x3C, 0x6D))
    add_textbox(slide11, 0.8, 1.8, 10, 4.5, (
        "Age Distribution of Respondents:\n\n"
        "18-24: 14.2% (686 respondents)\n"
        "25-34: 28.7% (1,387 respondents)\n"
        "35-44: 24.1% (1,165 respondents)\n"
        "45-54: 18.3% (884 respondents)\n"
        "55-64: 10.4% (502 respondents)\n"
        "65+: 4.3% (208 respondents)\n\n"
        "Gender: Female 51.2%, Male 47.1%, Non-binary 1.7%"
    ), font_size=12, color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 12: Appendix - Technology Stack Analysis ---
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide12, 0.8, 0.5, 10, 1, "Appendix — Technology Stack",
                font_size=28, bold=True, color=RGBColor(0x1A, 0x3C, 0x6D))
    add_textbox(slide12, 0.8, 1.8, 10, 4.5, (
        "Primary Platforms Used by Respondents:\n\n"
        "Mobile (iOS): 38.4%\n"
        "Mobile (Android): 31.2%\n"
        "Desktop (Windows): 16.8%\n"
        "Desktop (macOS): 8.9%\n"
        "Tablet: 3.1%\n"
        "Other: 1.6%\n\n"
        "Browser Preference: Chrome 64%, Safari 18%, Firefox 9%, Edge 7%, Other 2%"
    ), font_size=12, color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 13: Recommendations ---
    slide13 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide13, 0.8, 0.5, 10, 1, "Recommendations",
                font_size=28, bold=True, color=RGBColor(0x1A, 0x3C, 0x6D))
    add_textbox(slide13, 0.8, 1.8, 10, 4, (
        "Based on our findings, we recommend the following actions:\n\n"
        "1. Accelerate mobile platform investment — allocate 40% of dev budget\n"
        "2. Expand Asia-Pacific presence with localized marketing campaigns\n"
        "3. Introduce tiered subscription model by Q2 2026\n"
        "4. Enhance personalization engine using ML-driven recommendations\n"
        "5. Establish sustainability reporting framework for ESG compliance\n"
        "6. Pilot voice commerce integration by Q3 2026"
    ), font_size=14, color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 14: Next Steps ---
    slide14 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide14, 0.8, 0.5, 10, 1, "Next Steps",
                font_size=28, bold=True, color=RGBColor(0x1A, 0x3C, 0x6D))
    add_textbox(slide14, 0.8, 1.8, 10, 4, (
        "Timeline for Implementation:\n\n"
        "January 2026: Finalize strategy with executive leadership\n"
        "February 2026: Begin mobile platform redesign sprint\n"
        "March 2026: Launch Asia-Pacific market research phase\n"
        "April 2026: Subscription model beta testing\n"
        "June 2026: Full rollout of personalization features\n"
        "September 2026: Voice commerce pilot launch\n\n"
        "Quarterly reviews to assess progress against KPIs."
    ), font_size=14, color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 15: Thank You ---
    slide15 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide15, 1.5, 2.5, 10, 1.5, "Thank You",
                font_size=36, bold=True, alignment=PP_ALIGN.CENTER,
                color=RGBColor(0x1A, 0x3C, 0x6D))
    add_textbox(slide15, 1.5, 4.2, 10, 1, (
        "Questions? Contact: elena.vasquez@company.com\n"
        "Research Division | Strategic Insights Team"
    ), font_size=16, alignment=PP_ALIGN.CENTER,
       color=RGBColor(0x70, 0x70, 0x70))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
