"""
FINAL REWARD SCRIPT - SUCCESS
Task: Quick LibreOffice Impress question: on slide 124 I'd like the title to read exactly "Method" and have the Outline text effect (the hollow lettering). What's the fastest way to set that up?
Generated: 2025-09-10 23:13:58
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation


def _has_outline_text_effect(shape):
    """Detects the "Outline" (hollow) text effect on a shape.

    The typical XML pattern for outline (hollow) text in PPTX is:
        <a:rPr>
            <a:noFill/>        <-- no interior fill
            <a:ln w="...">    <-- stroke/outline definition
                 ...
            </a:ln>
        </a:rPr>
    We therefore look for BOTH a:noFill and a:ln tags under the run properties.
    """
    has_no_fill = False
    has_line = False

    # Traverse the underlying XML of the shape to look for the tags.
    for elem in shape._element.iter():
        tag = elem.tag
        if tag.endswith('}noFill'):
            has_no_fill = True
        elif tag.endswith('}ln'):
            has_line = True

        # Early-exit if both conditions are already met
        if has_no_fill and has_line:
            return True

    return False


def verify_task(file_path):
    """Verify the Impress / PowerPoint task.

    Task requirements:
      1. Slide 124 must exist.
      2. The title on slide 124 must read exactly "Method" (case-sensitive).
      3. That title text must use the Outline (hollow lettering) effect.

    Progressive scoring:
      • 0.4 points for finding the exact title text.
      • 0.6 additional points for detecting the outline effect.
      • Score is capped at 1.0.
    """
    print(f"Verifying presentation: {file_path}")

    if not os.path.exists(file_path):
        print("✗ File does not exist")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Requirement: slide 124 must exist (index 123)
    if len(prs.slides) < 124:
        print(f"✗ Slide 124 missing – only {len(prs.slides)} slides found")
        print("REWARD: 0.0")
        return 0.0

    print("✓ Slide 124 exists")
    slide124 = prs.slides[123]  # zero-based index

    title_found = False
    outline_ok = False

    # Search for the exact title text "Method"
    for shape in slide124.shapes:
        if hasattr(shape, "text"):
            if shape.text.strip() == "Method":
                title_found = True
                outline_ok = _has_outline_text_effect(shape)
                break

    score = 0.0

    if title_found:
        print("✓ Title text \"Method\" found")
        score += 0.4
        if outline_ok:
            print("✓ Outline text effect detected")
            score += 0.6
        else:
            print("✗ Outline text effect NOT detected")
    else:
        print("✗ Title text \"Method\" not found on slide 124")

    final_score = min(score, 1.0)
    print(f"REWARD: {final_score}")
    return final_score


if __name__ == "__main__":
    FILE = "/home/user/quick_libreoffice_impress_question_on_slide_124_id_like_the_title_to_read_exactly_method_and_have_th_golden.pptx"
    verify_task(FILE)
