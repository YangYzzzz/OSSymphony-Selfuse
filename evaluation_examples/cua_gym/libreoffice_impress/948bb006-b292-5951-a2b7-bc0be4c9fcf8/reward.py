"""
Reward Script: Waterfall chart on slide 5 of Cash_Flow.pptx
Task ID: impress_gf2_038
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Slide 5 contains a stacked bar chart
  Component 2 (0.25): Chart has 7 correct categories for cash flow waterfall
  Component 3 (0.25): Chart has 2 series with correct waterfall data values
  Component 4 (0.25): Chart uses green (#16A34A) for positive and red (#DC2626) for negative bars
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_038'

# Expected categories (normalized — stripped and lowered for flexible matching)
EXPECTED_CATEGORIES = [
    'starting balance',
    'revenue',
    'cogs',
    'opex',
    'tax',
    'net cash',
    'ending balance',
]

# Expected visible-bar values (Series 1 in waterfall stacked bar)
EXPECTED_VALUES = [100.0, 75.0, 40.0, 25.0, 15.0, 35.0, 130.0]

# Expected base values (Series 0 — transparent bottom segments)
EXPECTED_BASE = [0.0, 100.0, 135.0, 110.0, 95.0, 95.0, 0.0]

# Color mapping: index -> expected color for the visible series data points
# Green for positive/total bars, red for negative bars
EXPECTED_COLORS = {
    0: '16A34A',  # Starting Balance (green — total)
    1: '16A34A',  # Revenue (green — positive)
    2: 'DC2626',  # COGS (red — negative)
    3: 'DC2626',  # OpEx (red — negative)
    4: 'DC2626',  # Tax (red — negative)
    5: '16A34A',  # Net Cash (green — positive)
    6: '16A34A',  # Ending Balance (green — total)
}


def normalize_cat(s):
    """Normalize category string for comparison."""
    return ' '.join(s.replace('\n', ' ').split()).strip().lower()


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide5 = prs.slides[4]

    # Find chart shape on slide 5
    chart_shape = None
    for shape in slide5.shapes:
        if shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Slide 5 contains a stacked bar chart (0.25 points)
    try:
        if chart_shape is None:
            print("FAIL: Component 1 — No chart found on slide 5")
        else:
            chart = chart_shape.chart
            chart_type = chart.chart_type
            # BAR_STACKED = 58, BAR_CLUSTERED = 57; accept stacked bar variants
            # python-pptx chart_type enum: BAR_STACKED=58, BAR_STACKED_100=59
            if chart_type in (57, 58, 59):
                print(f"PASS: Component 1 — Stacked/clustered bar chart found on slide 5, type={chart_type} (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 — Chart found but type={chart_type}, expected stacked bar (57/58/59)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    if chart_shape is None:
        # No chart means remaining components cannot pass
        print("REWARD: 0.0")
        return 0.0

    chart = chart_shape.chart
    plot = chart.plots[0]

    # Component 2: Chart has 7 correct categories (0.25 points)
    try:
        cats = [str(c) for c in plot.categories]
        norm_cats = [normalize_cat(c) for c in cats]
        if len(norm_cats) == 7:
            matched = sum(1 for nc, ec in zip(norm_cats, EXPECTED_CATEGORIES) if nc == ec)
            if matched >= 6:
                print(f"PASS: Component 2 — 7 categories found, {matched}/7 match expected names (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 2 — Categories found but only {matched}/7 match. Got: {norm_cats}")
        else:
            print(f"FAIL: Component 2 — Expected 7 categories, found {len(norm_cats)}: {norm_cats}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Chart has 2 series with correct waterfall data values (0.25 points)
    try:
        num_series = len(plot.series)
        if num_series < 2:
            print(f"FAIL: Component 3 — Expected 2 series for waterfall, found {num_series}")
        else:
            # Series 0 = base (transparent), Series 1 = visible bars
            base_vals = list(plot.series[0].values)
            vis_vals = list(plot.series[1].values)

            # Check visible values match expected
            vis_match = all(
                abs(a - b) < 1.0
                for a, b in zip(vis_vals, EXPECTED_VALUES)
            ) if len(vis_vals) == 7 else False

            # Check base values match expected
            base_match = all(
                abs(a - b) < 1.0
                for a, b in zip(base_vals, EXPECTED_BASE)
            ) if len(base_vals) == 7 else False

            if vis_match and base_match:
                print(f"PASS: Component 3 — Both series have correct waterfall values (0.25 pts)")
                print(f"  Base: {base_vals}")
                print(f"  Visible: {vis_vals}")
                total_score += 0.25
            elif vis_match:
                print(f"PARTIAL: Component 3 — Visible values correct but base values wrong (0.15 pts)")
                print(f"  Base expected: {EXPECTED_BASE}, got: {base_vals}")
                total_score += 0.15
            else:
                print(f"FAIL: Component 3 — Values don't match.")
                print(f"  Expected visible: {EXPECTED_VALUES}, got: {vis_vals}")
                print(f"  Expected base: {EXPECTED_BASE}, got: {base_vals}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Correct colors — green for positive, red for negative (0.25 points)
    try:
        ns = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
        }

        # Check Series 0 has noFill (transparent base)
        ser0_el = plot.series[0]._element
        has_nofill = ser0_el.find('.//a:noFill', ns) is not None

        # Check Series 1 data point colors
        ser1_el = plot.series[1]._element
        dpts = ser1_el.findall('.//c:dPt', ns)
        color_matches = 0
        color_total = len(EXPECTED_COLORS)

        for dpt in dpts:
            idx_el = dpt.find('c:idx', ns)
            if idx_el is None:
                continue
            idx_val = int(idx_el.get('val'))
            clr_el = dpt.find('.//a:solidFill/a:srgbClr', ns)
            if clr_el is not None and idx_val in EXPECTED_COLORS:
                actual_color = clr_el.get('val', '').upper()
                expected_color = EXPECTED_COLORS[idx_val].upper()
                if actual_color == expected_color:
                    color_matches += 1
                else:
                    print(f"  Point {idx_val}: expected {expected_color}, got {actual_color}")

        if has_nofill and color_matches >= 5:
            print(f"PASS: Component 4 — Base series transparent, {color_matches}/{color_total} point colors correct (0.25 pts)")
            total_score += 0.25
        elif color_matches >= 5:
            print(f"PARTIAL: Component 4 — {color_matches}/{color_total} colors correct but base not transparent (0.15 pts)")
            total_score += 0.15
        elif color_matches >= 3:
            print(f"PARTIAL: Component 4 — {color_matches}/{color_total} colors correct (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 4 — Only {color_matches}/{color_total} colors correct, base noFill={has_nofill}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist app state before verification (LibreOffice may have unsaved changes)
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
