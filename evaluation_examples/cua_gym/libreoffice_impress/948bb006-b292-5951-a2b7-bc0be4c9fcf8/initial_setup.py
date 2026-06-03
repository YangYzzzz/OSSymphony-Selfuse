"""
Initial Setup: Cash Flow Waterfall presentation with 8 slides, slide 5 empty for chart
Task ID: impress_gf2_038
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_038'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_slide(prs, layout_idx, title_text, body_lines):
    """Add a slide with title and bullet-point body text."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    # Find content placeholder (usually index 1)
    body_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            body_ph = ph
            break
    if body_ph and body_lines:
        tf = body_ph.text_frame
        tf.clear()
        for i, line in enumerate(body_lines):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                p = tf.add_paragraph()
                p.text = line
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Cash Flow Analysis Q4 2025"
    slide1.placeholders[1].text = "Prepared by Finance Department\nNovember 2025"

    # --- Slide 2: Executive Summary ---
    add_text_slide(prs, 1, "Executive Summary", [
        "Overall cash position improved by 30% compared to Q3",
        "Revenue growth driven by enterprise software segment",
        "Operating expenses reduced through process optimization",
        "Tax obligations met ahead of schedule",
        "Net cash flow positive for the third consecutive quarter",
    ])

    # --- Slide 3: Revenue Breakdown ---
    add_text_slide(prs, 1, "Revenue Breakdown", [
        "Enterprise Software Licenses: $42.3M (+18% YoY)",
        "Cloud Services Subscription: $21.7M (+35% YoY)",
        "Professional Services: $8.5M (+5% YoY)",
        "Maintenance & Support: $2.5M (-2% YoY)",
        "Total Revenue: $75M for the quarter",
    ])

    # --- Slide 4: Cost Structure ---
    add_text_slide(prs, 1, "Cost Structure Overview", [
        "Cost of Goods Sold (COGS): $40M - includes hosting, licenses",
        "Operating Expenses: $25M - salaries, marketing, R&D",
        "Tax Provisions: $15M - federal and state obligations",
        "Capital Expenditures: $3.2M - infrastructure upgrades",
        "Total outflows managed within budget targets",
    ])

    # --- Slide 5: Cash Flow Waterfall Analysis (EMPTY - no chart) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title text box manually
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Cash Flow Waterfall Analysis"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 6: Investment Outlook ---
    add_text_slide(prs, 1, "Investment Outlook", [
        "Planned R&D investment of $15M in Q1 2026",
        "New data center expansion in Southeast Asia",
        "Strategic acquisition pipeline under review",
        "Expected ROI of 22% on cloud infrastructure upgrades",
        "Board approved $8M for AI integration initiatives",
    ])

    # --- Slide 7: Risk Assessment ---
    add_text_slide(prs, 1, "Risk Assessment", [
        "Currency fluctuation exposure: Medium risk",
        "Supply chain dependencies: Low risk (diversified)",
        "Regulatory compliance: Monitoring EU data sovereignty rules",
        "Market competition: Aggressive pricing from new entrants",
        "Mitigation strategies in place for top 3 risks",
    ])

    # --- Slide 8: Summary & Next Steps ---
    add_text_slide(prs, 1, "Summary & Next Steps", [
        "Cash position remains strong at $130M ending balance",
        "Continue cost optimization program into Q1 2026",
        "Expand revenue diversification into Asian markets",
        "Schedule quarterly cash flow review for January 15",
        "Action items assigned to department heads by Dec 1",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
