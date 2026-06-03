"""
Initial Setup: Create interactive quiz presentation with 10 slides (no navigation buttons)
Task ID: impress_gf5_018
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_018'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_centered_title(slide, title_text, subtitle_text=None):
    """Add title and optional subtitle to a slide using text boxes."""
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x3A, 0x87)

    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = subtitle_text
        run2.font.size = Pt(20)
        run2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)


def add_content_text(slide, text, top=Inches(2.5)):
    """Add content text to a slide."""
    txBox = slide.shapes.add_textbox(Inches(1), top, Inches(8), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Menu ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only / Blank
    add_centered_title(slide1, "Interactive Science Quiz", "Select a question to begin")

    # --- Slide 2: Answer 1 ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_centered_title(slide2, "Answer 1")
    add_content_text(slide2,
        "The correct answer is B) 9.8 m/s².\n\n"
        "The acceleration due to gravity on Earth's surface is approximately "
        "9.8 meters per second squared. This value was first measured accurately "
        "by Galileo Galilei through his famous inclined plane experiments.")

    # --- Slide 3: Question 1 ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_centered_title(slide3, "Question 1: Physics")
    add_content_text(slide3,
        "What is the approximate acceleration due to gravity on Earth?\n\n"
        "A) 5.2 m/s²\n"
        "B) 9.8 m/s²\n"
        "C) 12.4 m/s²\n"
        "D) 15.0 m/s²")

    # --- Slide 4: Answer 2 ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_centered_title(slide4, "Answer 2")
    add_content_text(slide4,
        "The correct answer is C) H₂O.\n\n"
        "Water is composed of two hydrogen atoms bonded to one oxygen atom. "
        "This molecular structure gives water its unique properties, including "
        "its high specific heat capacity and solvent capabilities.")

    # --- Slide 5: Question 2 ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_centered_title(slide5, "Question 2: Chemistry")
    add_content_text(slide5,
        "What is the chemical formula for water?\n\n"
        "A) CO₂\n"
        "B) NaCl\n"
        "C) H₂O\n"
        "D) O₂")

    # --- Slide 6: Answer 3 ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_centered_title(slide6, "Answer 3")
    add_content_text(slide6,
        "The correct answer is A) Mitochondria.\n\n"
        "Often called the 'powerhouse of the cell,' mitochondria generate most "
        "of the cell's supply of adenosine triphosphate (ATP), used as a source "
        "of chemical energy through oxidative phosphorylation.")

    # --- Slide 7: Question 3 ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_centered_title(slide7, "Question 3: Biology")
    add_content_text(slide7,
        "Which organelle is known as the 'powerhouse of the cell'?\n\n"
        "A) Mitochondria\n"
        "B) Nucleus\n"
        "C) Ribosome\n"
        "D) Golgi apparatus")

    # --- Slide 8: Answer 4 ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_centered_title(slide8, "Answer 4")
    add_content_text(slide8,
        "The correct answer is D) Mars.\n\n"
        "Mars is called the 'Red Planet' because of the iron oxide (rust) "
        "prevalent on its surface, giving it a reddish appearance visible "
        "from Earth with the naked eye.")

    # --- Slide 9: Question 4 ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_centered_title(slide9, "Question 4: Astronomy")
    add_content_text(slide9,
        "Which planet is known as the 'Red Planet'?\n\n"
        "A) Venus\n"
        "B) Jupiter\n"
        "C) Saturn\n"
        "D) Mars")

    # --- Slide 10: Results ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_centered_title(slide10, "Quiz Results", "Review your answers and check your score")
    add_content_text(slide10,
        "Congratulations on completing the Interactive Science Quiz!\n\n"
        "Topics covered:\n"
        "• Physics - Gravity\n"
        "• Chemistry - Molecular formulas\n"
        "• Biology - Cell organelles\n"
        "• Astronomy - Planets",
        top=Inches(3.5))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
