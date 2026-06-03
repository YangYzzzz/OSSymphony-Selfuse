"""
Reward Script: Hyperlink navigation buttons in interactive_quiz.pptx
Task ID: impress_gf5_018
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Slide 1 has 4 button shapes labeled Q1, Q2, Q3, Q4
  Component 2 (0.30): Q1->slide3, Q2->slide5, Q3->slide7, Q4->slide9 via click actions
  Component 3 (0.20): Slides 3,5,7,9 each have a 'Back to Menu' button shape
  Component 4 (0.25): 'Back to Menu' buttons all link back to slide 1
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_018'


def persist_app_state(domain: str):
    """Save any unsaved GUI edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def find_target_slide_index(prs, target_slide):
    """Given a target_slide object from click_action, find its 0-based index in prs.slides."""
    for idx, s in enumerate(prs.slides):
        if s._element is target_slide._element:
            return idx
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 10 slides
    if len(prs.slides) < 10:
        print(f"PRECONDITION FAIL: Expected >= 10 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # ---- Component 1: Slide 1 has 4 button shapes labeled Q1, Q2, Q3, Q4 (0.25 pts) ----
    try:
        slide1 = prs.slides[0]
        # Find AUTO_SHAPE (type 1) shapes whose text matches Q1-Q4
        expected_labels = {"Q1", "Q2", "Q3", "Q4"}
        found_labels = set()
        for shape in slide1.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                text = shape.text.strip() if hasattr(shape, 'text') else ""
                if text in expected_labels:
                    found_labels.add(text)

        if found_labels == expected_labels:
            print(f"PASS: Component 1 — Slide 1 has all 4 button shapes: {sorted(found_labels)} (0.25 pts)")
            total_score += 0.25
        else:
            missing = expected_labels - found_labels
            print(f"FAIL: Component 1 — Missing button shapes on slide 1: {sorted(missing)} (found: {sorted(found_labels)})")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # ---- Component 2: Q buttons link to correct target slides (0.30 pts) ----
    # Q1->slide3 (idx 2), Q2->slide5 (idx 4), Q3->slide7 (idx 6), Q4->slide9 (idx 8)
    try:
        expected_targets = {"Q1": 2, "Q2": 4, "Q3": 6, "Q4": 8}
        correct_links = 0
        total_links = 4

        for shape in slide1.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                text = shape.text.strip() if hasattr(shape, 'text') else ""
                if text in expected_targets:
                    ca = shape.click_action
                    if ca is not None and ca.target_slide is not None:
                        target_idx = find_target_slide_index(prs, ca.target_slide)
                        expected_idx = expected_targets[text]
                        if target_idx == expected_idx:
                            correct_links += 1
                            print(f"  PASS: {text} -> Slide {target_idx + 1} (correct)")
                        else:
                            print(f"  FAIL: {text} -> Slide {target_idx + 1 if target_idx is not None else 'None'} (expected Slide {expected_idx + 1})")
                    else:
                        print(f"  FAIL: {text} has no click_action or no target_slide")

        if correct_links == total_links:
            print(f"PASS: Component 2 — All 4 Q buttons link to correct slides (0.30 pts)")
            total_score += 0.30
        elif correct_links > 0:
            partial = 0.30 * (correct_links / total_links)
            print(f"PARTIAL: Component 2 — {correct_links}/{total_links} correct links ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No Q buttons have correct hyperlinks")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # ---- Component 3: Slides 3,5,7,9 each have a 'Back to Menu' button (0.20 pts) ----
    try:
        back_button_slides = [2, 4, 6, 8]  # 0-based indices for slides 3,5,7,9
        found_back_buttons = 0

        for si in back_button_slides:
            slide = prs.slides[si]
            has_back = False
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    text = shape.text.strip().lower() if hasattr(shape, 'text') else ""
                    if "back" in text and "menu" in text:
                        has_back = True
                        break
            if has_back:
                found_back_buttons += 1
                print(f"  PASS: Slide {si + 1} has 'Back to Menu' button")
            else:
                print(f"  FAIL: Slide {si + 1} missing 'Back to Menu' button")

        if found_back_buttons == 4:
            print(f"PASS: Component 3 — All 4 question slides have 'Back to Menu' buttons (0.20 pts)")
            total_score += 0.20
        elif found_back_buttons > 0:
            partial = 0.20 * (found_back_buttons / 4)
            print(f"PARTIAL: Component 3 — {found_back_buttons}/4 slides have 'Back to Menu' ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No question slides have 'Back to Menu' buttons")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # ---- Component 4: 'Back to Menu' buttons link back to slide 1 (0.25 pts) ----
    try:
        correct_back_links = 0

        for si in back_button_slides:
            slide = prs.slides[si]
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    text = shape.text.strip().lower() if hasattr(shape, 'text') else ""
                    if "back" in text and "menu" in text:
                        ca = shape.click_action
                        if ca is not None and ca.target_slide is not None:
                            target_idx = find_target_slide_index(prs, ca.target_slide)
                            if target_idx == 0:  # slide 1
                                correct_back_links += 1
                                print(f"  PASS: Slide {si + 1} 'Back to Menu' -> Slide 1 (correct)")
                            else:
                                print(f"  FAIL: Slide {si + 1} 'Back to Menu' -> Slide {target_idx + 1 if target_idx is not None else 'None'} (expected Slide 1)")
                        else:
                            print(f"  FAIL: Slide {si + 1} 'Back to Menu' has no click_action or no target")
                        break  # only check first matching button per slide

        if correct_back_links == 4:
            print(f"PASS: Component 4 — All 'Back to Menu' buttons link to Slide 1 (0.25 pts)")
            total_score += 0.25
        elif correct_back_links > 0:
            partial = 0.25 * (correct_back_links / 4)
            print(f"PARTIAL: Component 4 — {correct_back_links}/4 back links correct ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — No 'Back to Menu' buttons link to Slide 1")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point — persist then verify
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
