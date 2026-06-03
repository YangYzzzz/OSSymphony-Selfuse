"""
Reward Script: Insert disclaimer text box on slide 2
Task ID: impress_gf3_004
Domain: libreoffice_impress
Scoring:
  Component 1 (0.35): Text box on slide 2 with exact text 'Beta Version - Do Not Distribute'
  Component 2 (0.25): Font is 10pt and italic
  Component 3 (0.20): Font color is #CC0000
  Component 4 (0.20): Text box positioned in bottom-left area of slide
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_004'


def persist_app_state(domain: str):
    """Save any unsaved LibreOffice changes before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(1.0)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def find_disclaimer_textbox(slide):
    """
    Search all shapes on a slide for a TEXT_BOX containing the disclaimer text.
    Returns (shape, run) if found, else (None, None).
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    target_text = 'Beta Version - Do Not Distribute'
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shape.has_text_frame:
            full_text = shape.text_frame.text.strip()
            if full_text == target_text:
                # Get first non-empty run
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if (run.text or "").strip():
                            return shape, run
    return None, None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Pt, Emu
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 2")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]  # 0-indexed, slide 2
    shape, run = find_disclaimer_textbox(slide2)

    # Component 1: Text box on slide 2 with exact text (0.35 points)
    try:
        if shape is not None and run is not None:
            print(f"PASS: Component 1 - Text box found with text 'Beta Version - Do Not Distribute' (0.35 pts)")
            total_score += 0.35
        else:
            print(f"FAIL: Component 1 - No text box with exact text 'Beta Version - Do Not Distribute' found on slide 2")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Font is 10pt and italic (0.25 points)
    try:
        if run is not None:
            font_size = run.font.size
            is_italic = run.font.italic
            # Normalize None -> False for italic
            is_italic_actual = True if is_italic is True else False
            # 10pt = Pt(10) = 127000 EMU
            size_ok = (font_size is not None and abs(font_size - 127000) < 5000)
            italic_ok = is_italic_actual

            if size_ok and italic_ok:
                print(f"PASS: Component 2 - Font is 10pt (size={font_size} EMU) and italic={is_italic} (0.25 pts)")
                total_score += 0.25
            elif size_ok:
                print(f"PARTIAL: Component 2 - Size correct ({font_size} EMU) but italic={is_italic} (0.125 pts)")
                total_score += 0.125
            elif italic_ok:
                print(f"PARTIAL: Component 2 - Italic correct but size={font_size} EMU, expected ~127000 (0.125 pts)")
                total_score += 0.125
            else:
                print(f"FAIL: Component 2 - size={font_size} EMU (expected ~127000), italic={is_italic} (expected True)")
        else:
            print(f"FAIL: Component 2 - No disclaimer text box found, cannot check font")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Font color is #CC0000 (0.20 points)
    try:
        if run is not None:
            try:
                if run.font.color.type is not None:
                    actual_rgb = str(run.font.color.rgb).upper()
                    if actual_rgb == 'CC0000':
                        print(f"PASS: Component 3 - Font color is #CC0000 (0.20 pts)")
                        total_score += 0.20
                    else:
                        print(f"FAIL: Component 3 - Font color is #{actual_rgb}, expected #CC0000")
                else:
                    print(f"FAIL: Component 3 - Font color type is None (no explicit color set)")
            except AttributeError:
                print(f"FAIL: Component 3 - Font color is theme-based or not accessible")
        else:
            print(f"FAIL: Component 3 - No disclaimer text box found, cannot check color")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Position in bottom-left area (0.20 points)
    # Slide is 10in wide (9144000 EMU) x 7.5in tall (6858000 EMU)
    # Bottom-left means: left < 25% of slide width AND top > 75% of slide height
    # Context says approximately x=0.3cm, y=17.5cm on a 25.4cm tall slide
    # In standard slide: top should be > ~80% of height, left should be < ~20% of width
    try:
        if shape is not None:
            slide_width = prs.slide_width   # 9144000
            slide_height = prs.slide_height  # 6858000

            left_ratio = shape.left / slide_width
            top_ratio = shape.top / slide_height

            # Bottom-left: left should be small (<30% of width), top should be large (>75% of height)
            left_ok = left_ratio < 0.30
            top_ok = top_ratio > 0.75

            print(f"  DEBUG: Position - left={shape.left} ({left_ratio:.2%} of width), top={shape.top} ({top_ratio:.2%} of height)")

            if left_ok and top_ok:
                print(f"PASS: Component 4 - Text box is in bottom-left area (0.20 pts)")
                total_score += 0.20
            elif top_ok:
                print(f"PARTIAL: Component 4 - Bottom position correct but left_ratio={left_ratio:.2%} (expected <30%) (0.10 pts)")
                total_score += 0.10
            elif left_ok:
                print(f"PARTIAL: Component 4 - Left position correct but top_ratio={top_ratio:.2%} (expected >75%) (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 4 - Position not bottom-left: left_ratio={left_ratio:.2%}, top_ratio={top_ratio:.2%}")
        else:
            print(f"FAIL: Component 4 - No disclaimer text box found, cannot check position")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist unsaved state before checking
persist_app_state("libreoffice_impress")

# Test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
