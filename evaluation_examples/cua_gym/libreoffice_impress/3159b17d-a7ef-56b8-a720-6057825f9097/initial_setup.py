"""
Initial Setup: Insert disclaimer text box on slide 2
Task ID: impress_gf3_004
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_004'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard 10x7.5 widescreen
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Software Preview"
    slide1.placeholders[1].text = "Internal Product Review — Q2 2025"

    # --- Slide 2: Feature Overview (target slide) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Feature Overview"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Real-time collaboration with up to 50 concurrent users"
    p2a = tf2.add_paragraph()
    p2a.text = "Advanced analytics dashboard with customizable widgets"
    p2a.level = 0
    p2b = tf2.add_paragraph()
    p2b.text = "End-to-end encryption for all data in transit and at rest"
    p2b.level = 0
    p2c = tf2.add_paragraph()
    p2c.text = "REST API with OAuth 2.0 authentication support"
    p2c.level = 0
    p2d = tf2.add_paragraph()
    p2d.text = "Automated backup and disaster recovery system"
    p2d.level = 0

    # --- Slide 3: Architecture ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "System Architecture"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Microservices-based backend deployed on Kubernetes"
    p3a = tf3.add_paragraph()
    p3a.text = "PostgreSQL primary database with Redis caching layer"
    p3a.level = 0
    p3b = tf3.add_paragraph()
    p3b.text = "Message queue powered by Apache Kafka for async processing"
    p3b.level = 0
    p3c = tf3.add_paragraph()
    p3c.text = "CDN-accelerated static asset delivery via CloudFront"
    p3c.level = 0

    # --- Slide 4: Performance Metrics ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Performance Metrics"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Average API response time: 45ms (p99: 120ms)"
    p4a = tf4.add_paragraph()
    p4a.text = "System uptime: 99.97% over the last 12 months"
    p4a.level = 0
    p4b = tf4.add_paragraph()
    p4b.text = "Peak throughput: 12,000 requests per second"
    p4b.level = 0
    p4c = tf4.add_paragraph()
    p4c.text = "Database query optimization reduced latency by 38%"
    p4c.level = 0

    # --- Slide 5: Roadmap ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Product Roadmap"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Q3 2025: Mobile app launch (iOS and Android)"
    p5a = tf5.add_paragraph()
    p5a.text = "Q4 2025: Multi-language support (12 languages)"
    p5a.level = 0
    p5b = tf5.add_paragraph()
    p5b.text = "Q1 2026: Enterprise SSO integration"
    p5b.level = 0
    p5c = tf5.add_paragraph()
    p5c.text = "Q2 2026: AI-powered workflow automation"
    p5c.level = 0

    # --- Slide 6: Questions & Contact ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Questions & Contact"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Product Lead: Rachel Kim — rachel.kim@techcorp.io"
    p6a = tf6.add_paragraph()
    p6a.text = "Engineering: dev-team@techcorp.io"
    p6a.level = 0
    p6b = tf6.add_paragraph()
    p6b.text = "Documentation: https://docs.internal.techcorp.io/preview"
    p6b.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
