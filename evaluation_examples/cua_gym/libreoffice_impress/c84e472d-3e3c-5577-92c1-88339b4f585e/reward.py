"""
FINAL REWARD SCRIPT - SUCCESS
Task: Shade every alternate row in Table 1 with 10% Gray (banded rows).
Generated: 2025-10-17 12:32:31
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation
from pptx.dml.color import RGBColor

"""
Reward Script: Verify that every alternate row in "Table 1" is shaded with 10% Gray
("banded rows") in the supplied PPTX file and award a progressive score.

Scoring rubric (max 1.0):
• 0.4 – Required table (named "Table 1" – or the first table if name missing) exists
• 0.3 – Rows exhibit a strict alternate shading pattern (banded rows)
• 0.3 – Shaded rows use a colour that is ~10 % gray (RGB≈230,230,230) within tolerance

The script prints detailed diagnostics and the final reward in the mandatory
format:  "REWARD: X.X"  (float between 0-1).
"""

FILE_PATH = "/home/user/shade_every_alternate_row_in_table_1_with_10_gray_banded_rows.pptx"

# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------

def is_gray(rgb_tuple, tol=15):
    """Return True if the colour is approximately gray (all channels similar)."""
    if rgb_tuple is None:
        return False
    r, g, b = rgb_tuple
    return abs(r - g) <= tol and abs(r - b) <= tol and abs(g - b) <= tol


def is_ten_percent_gray(rgb_tuple, tol_channel=25):
    """10 % gray ≈ RGB(230,230,230). Accept a reasonable tolerance."""
    if rgb_tuple is None or not is_gray(rgb_tuple):
        return False
    r, g, b = rgb_tuple
    return 210 <= r <= 245  # accept slight deviations around 230


def extract_rgb(pptx_rgb):
    """Convert python-pptx RGBColor to (r,g,b) tuple or None."""
    if pptx_rgb is None:
        return None
    if isinstance(pptx_rgb, RGBColor):
        return (pptx_rgb[0], pptx_rgb[1], pptx_rgb[2])
    # Fallback for hex-string instances
    try:
        return tuple(pptx_rgb)
    except Exception:
        return None


# ---------------------------------------------------------------------------
# Core verification logic
# ---------------------------------------------------------------------------

def verify_alternate_gray_shading(table):
    """Return partial score (0-0.6) based on alternation & correct 10 % gray colour."""
    shaded_pattern_ok = True
    gray_colour_ok = True

    shaded_flags = []  # True if row appears shaded (solid fill applied)

    for idx, row in enumerate(table.rows):
        cell = row.cells[0]  # inspect first cell – shading is row-wide
        fill = cell.fill
        rgb_tuple = None
        shaded = False

        if fill.type == 1:  # solid fill present
            rgb_tuple = extract_rgb(fill.fore_color.rgb)
            shaded = rgb_tuple is not None
        shaded_flags.append(shaded)

        # If this row is shaded, validate it is 10 % gray
        if shaded and not is_ten_percent_gray(rgb_tuple):
            gray_colour_ok = False

    # Verify true alternation pattern (no two consecutive same shading flag)
    if len(shaded_flags) < 2:
        shaded_pattern_ok = False  # need at least 2 rows to demonstrate alternation
    else:
        shaded_pattern_ok = all(shaded_flags[i] != shaded_flags[i + 1] for i in range(len(shaded_flags) - 1))

    # Scoring
    score = 0.0
    if shaded_pattern_ok:
        print("✓ Alternate banded pattern verified (0.3 points)")
        score += 0.3
    else:
        print("✗ Rows are NOT in strict alternate pattern (0 points)")

    if gray_colour_ok and any(shaded_flags):
        print("✓ Shaded rows use ~10 % gray (0.3 points)")
        score += 0.3
    else:
        print("✗ Shaded rows are not correct 10 % gray (0 points)")

    return score


def verify_task(file_path=FILE_PATH):
    total_score = 0.0

    # Basic existence (no points – prerequisite only)
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    # Load presentation
    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded PPTX successfully – slides: {len(prs.slides)}")
    except Exception as e:
        print(f"✗ Could not load PPTX: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Locate the required table
    target_table = None
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                if getattr(shape, "name", "") == "Table 1":
                    target_table = shape.table
                    break
        if target_table:
            break

    # Fallback: use first table if specifically named one not found
    if target_table is None:
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    target_table = shape.table
                    break
            if target_table:
                break

    # Verify table presence
    if target_table is None:
        print("✗ No table found – cannot verify task")
        print("REWARD: 0.0")
        return 0.0

    print("✓ Located table for verification (0.4 points)")
    total_score += 0.4

    # Verify banded shading & colour correctness
    total_score += verify_alternate_gray_shading(target_table)

    # Cap at 1.0
    final_score = min(total_score, 1.0)
    print(f"REWARD: {final_score}")
    return final_score


# ---------------------------------------------------------------------------
# Execute verification when run as main module
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    verify_task()

