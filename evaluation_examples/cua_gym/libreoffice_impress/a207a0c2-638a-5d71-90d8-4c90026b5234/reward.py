"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 109’s layout looks off—Table 1 isn’t where it should be. Please drop that table so its upper-left corner sits exactly at X = 2.0 cm and Y = 10.0 cm.
Generated: 2025-09-10 22:51:47
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
import math
from pptx import Presentation

"""
Reward script for the task:
"Slide 109’s layout looks off—Table 1 isn’t where it should be. Please drop that table so its upper-left corner sits exactly at X = 2.0 cm and Y = 10.0 cm."

The script verifies:
1. A PPTX file exists in /home/user (the working directory for these tasks).
2. The presentation contains at least 109 slides.
3. Slide 109 (index 108) contains **a table**.
4. The table’s upper-left corner is positioned at **exactly** X = 2.0 cm and Y = 10.0 cm (within a tight tolerance).

Progressive scoring (0-1 range):
• 0.0  – major requirements missing (no slide 109 or no table)
• +0.2 – table detected on slide 109
• +0.8 – table perfectly positioned (≤0.05 cm error)
       – smaller bonuses if within looser tolerances (0.6/0.4/0.2)

The script prints detailed diagnostics and finally prints
    REWARD: <score>
exactly as required.
"""

# Constants
EMU_PER_CM = 360000  # EMUs per centimetre in PPTX coordinate system
EXPECTED_LEFT_CM = 2.0
EXPECTED_TOP_CM = 10.0
EXPECTED_LEFT_EMU = EXPECTED_LEFT_CM * EMU_PER_CM
EXPECTED_TOP_EMU = EXPECTED_TOP_CM * EMU_PER_CM

# Locate a single pptx file in /home/user (task environment)
SEARCH_DIR = "/home/user"
pptx_files = [f for f in os.listdir(SEARCH_DIR) if f.lower().endswith(".pptx")]

pptx_path = None
if len(pptx_files) == 1:
    pptx_path = os.path.join(SEARCH_DIR, pptx_files[0])
else:
    # Prefer the file that does *not* contain the word "golden" if multiple
    non_golden = [f for f in pptx_files if "golden" not in f.lower()]
    pptx_path = os.path.join(SEARCH_DIR, (non_golden[0] if non_golden else pptx_files[0])) if pptx_files else None

if not pptx_path or not os.path.exists(pptx_path):
    print("✗ No PPTX file found for verification in", SEARCH_DIR)
    print("REWARD: 0.0")
    raise SystemExit

print("Verifying file:", pptx_path)

# Load presentation
try:
    prs = Presentation(pptx_path)
except Exception as e:
    print("✗ Failed to load presentation:", e)
    print("REWARD: 0.0")
    raise SystemExit

# Check slide count
if len(prs.slides) < 109:
    print(f"✗ Presentation has only {len(prs.slides)} slides; slide 109 is missing.")
    print("REWARD: 0.0")
    raise SystemExit

print("✓ Slide 109 found (index 108).")
slide_109 = prs.slides[108]

# Locate the table closest to the expected position
closest_table = None
closest_distance_cm = None
for shape in slide_109.shapes:
    if shape.has_table:
        dist_cm = math.hypot((shape.left - EXPECTED_LEFT_EMU) / EMU_PER_CM,
                             (shape.top  - EXPECTED_TOP_EMU)  / EMU_PER_CM)
        if closest_distance_cm is None or dist_cm < closest_distance_cm:
            closest_distance_cm = dist_cm
            closest_table = shape

score = 0.0
max_score = 1.0

if closest_table is None:
    print("✗ No table found on slide 109.")
else:
    print("✓ Table detected on slide 109.")
    score += 0.2  # credit for finding a table we can evaluate

    # Calculate position errors in cm
    left_cm = closest_table.left / EMU_PER_CM
    top_cm  = closest_table.top  / EMU_PER_CM
    h_err = abs(left_cm - EXPECTED_LEFT_CM)
    v_err = abs(top_cm  - EXPECTED_TOP_CM)

    print(f"Table position: left={left_cm:.2f} cm, top={top_cm:.2f} cm")
    print(f"Errors: horizontal={h_err:.2f} cm, vertical={v_err:.2f} cm")

    # Progressive accuracy scoring
    if h_err <= 0.05 and v_err <= 0.05:
        score += 0.8
        print("✓ Table position perfectly aligned (≤0.05 cm).")
    elif h_err <= 0.20 and v_err <= 0.20:
        score += 0.6
        print("✓ Table position very close (≤0.20 cm).")
    elif h_err <= 0.50 and v_err <= 0.50:
        score += 0.4
        print("✓ Table position close (≤0.50 cm).")
    elif h_err <= 1.00 and v_err <= 1.00:
        score += 0.2
        print("✓ Table position somewhat close (≤1.00 cm).")
    else:
        print("✗ Table far from expected position (>1 cm).")

# Cap score at 1.0
final_score = min(score, max_score)
print(f"Final score: {final_score}")
print(f"REWARD: {final_score}")
