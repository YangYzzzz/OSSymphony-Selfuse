"""
Initial Setup: Create a 6-slide training deck where slide 2 advances on mouse click only.
Task ID: impress_tm_020
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import zipfile
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_020'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_shape(shape, text, font_name="Arial", font_size=Pt(18),
                       bold=False, color=None, alignment=None):
    """Helper to set text on a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    for run in p.runs:
        run.font.name = font_name
        run.font.size = font_size
        run.font.bold = bold
        if color:
            run.font.color.rgb = color


def add_bullet_paragraph(text_frame, text, level=0, font_name="Arial",
                          font_size=Pt(16), color=None):
    """Add a bulleted paragraph to a text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    for run in p.runs:
        run.font.name = font_name
        run.font.size = font_size
        if color:
            run.font.color.rgb = color
    return p


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q2 2025 Sales Training"
    slide1.placeholders[1].text = "Accelerating Growth Through Customer-Centric Strategies"

    # ---- Slide 2: Agenda (this is the target slide) ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Training Agenda"
    body2 = slide2.placeholders[1]
    tf2 = body2.text_frame
    tf2.text = "Market Overview & Q1 Performance Review"
    add_bullet_paragraph(tf2, "New Product Line Launch Strategy", level=0)
    add_bullet_paragraph(tf2, "Target customer segments and verticals", level=1)
    add_bullet_paragraph(tf2, "Competitive Positioning & Objection Handling", level=0)
    add_bullet_paragraph(tf2, "Updated CRM Workflow & Pipeline Management", level=0)
    add_bullet_paragraph(tf2, "Role-Play Exercises & Q&A Session", level=0)

    # ---- Slide 3: Market Overview ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Market Overview"
    body3 = slide3.placeholders[1]
    tf3 = body3.text_frame
    tf3.text = "Total Addressable Market grew 12% YoY to $4.8B"
    add_bullet_paragraph(tf3, "Enterprise segment: $2.1B (+18%)", level=0)
    add_bullet_paragraph(tf3, "Mid-market segment: $1.6B (+9%)", level=0)
    add_bullet_paragraph(tf3, "SMB segment: $1.1B (+7%)", level=0)
    add_bullet_paragraph(tf3, "Key growth drivers: cloud migration, AI adoption, compliance requirements", level=0)

    # ---- Slide 4: Q1 Performance ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Q1 Performance Highlights"
    body4 = slide4.placeholders[1]
    tf4 = body4.text_frame
    tf4.text = "Revenue: $12.4M (106% of target)"
    add_bullet_paragraph(tf4, "New logos: 47 accounts (+23% vs Q4)", level=0)
    add_bullet_paragraph(tf4, "Average deal size: $263K (+11%)", level=0)
    add_bullet_paragraph(tf4, "Win rate: 34% (up from 29%)", level=0)
    add_bullet_paragraph(tf4, "Top performers: Sarah Chen ($2.1M), Marcus Rivera ($1.8M), Priya Sharma ($1.6M)", level=0)

    # ---- Slide 5: Competitive Positioning ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Competitive Landscape"
    body5 = slide5.placeholders[1]
    tf5 = body5.text_frame
    tf5.text = "Key competitors and our differentiators"
    add_bullet_paragraph(tf5, "Acme Corp: Strong in enterprise, weak in UX - emphasize our platform ease of use", level=0)
    add_bullet_paragraph(tf5, "NovaTech: Price leader - counter with TCO analysis and ROI data", level=0)
    add_bullet_paragraph(tf5, "CloudFirst: Good integrations - highlight our 200+ native connectors", level=0)
    add_bullet_paragraph(tf5, "Our edge: Customer success team, 99.9% uptime SLA, AI-powered analytics", level=0)

    # ---- Slide 6: Next Steps ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Action Items & Next Steps"
    body6 = slide6.placeholders[1]
    tf6 = body6.text_frame
    tf6.text = "Complete CRM migration by April 15th"
    add_bullet_paragraph(tf6, "Schedule practice sessions with regional managers", level=0)
    add_bullet_paragraph(tf6, "Review updated pricing sheets and proposal templates", level=0)
    add_bullet_paragraph(tf6, "Submit Q2 territory plans by end of week", level=0)
    add_bullet_paragraph(tf6, "Next training: Advanced Negotiation Techniques - May 8th", level=0)

    prs.save(OUTPUT)

    # Now ensure slide 2 has only advanceOnClick=true and NO auto-advance timing
    # Default python-pptx behavior is advanceOnClick, but let's be explicit via XML
    _ensure_slide2_click_only(OUTPUT)

    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


def _ensure_slide2_click_only(pptx_path):
    """Ensure slide 2 has advanceOnClick='1' and no advAfter timing."""
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    tmp_path = pptx_path + '.tmp'

    with zipfile.ZipFile(pptx_path, 'r') as zin:
        with zipfile.ZipFile(tmp_path, 'w') as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == 'ppt/slides/slide2.xml':
                    root = ET.fromstring(data)
                    # Remove any existing transition element
                    for tr in root.findall(f'{{{ns_p}}}transition'):
                        root.remove(tr)
                    # Add explicit transition with advClick only
                    # Insert transition element before the last element
                    tr_elem = ET.SubElement(root, f'{{{ns_p}}}transition')
                    tr_elem.set('advClick', '1')
                    # No advTm attribute = no auto-advance
                    data = ET.tostring(root, xml_declaration=True, encoding='UTF-8')
                zout.writestr(item, data)

    os.replace(tmp_path, pptx_path)


create_initial()
