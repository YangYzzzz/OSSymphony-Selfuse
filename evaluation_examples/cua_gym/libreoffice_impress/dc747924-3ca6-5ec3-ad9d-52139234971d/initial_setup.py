"""
Initial Setup: Master slide title placeholder at wrong position
Task ID: impress_fix_050
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_050'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard 10x7.5 slide dimensions
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Modify the slide master title placeholder to be at the WRONG position ---
    # Access the first slide master
    master = prs.slide_masters[0]

    # Find the title placeholder in the slide master and reposition it
    for ph in master.placeholders:
        if ph.placeholder_format.idx == 0:  # title placeholder
            ph.left = Inches(0.5)
            ph.top = Inches(3.0)
            ph.width = Inches(4.0)
            ph.height = Inches(1.25)
            break

    # Also fix it in all slide layouts that have a title placeholder
    for layout in prs.slide_layouts:
        for ph in layout.placeholders:
            if ph.placeholder_format.idx == 0:  # title placeholder
                ph.left = Inches(0.5)
                ph.top = Inches(3.0)
                ph.width = Inches(4.0)
                ph.height = Inches(1.25)

    # --- Create 10 slides with realistic content ---
    slide_contents = [
        {
            'layout': 0,  # Title Slide
            'title': 'Q3 2025 Strategic Review',
            'subtitle': 'Prepared by the Corporate Strategy Division\nSeptember 2025',
        },
        {
            'layout': 1,  # Title + Content
            'title': 'Executive Summary',
            'body': 'Revenue grew 18% year-over-year to $142M\nOperating margin improved to 23.5%\nNew product launches exceeded targets by 12%\nCustomer retention rate at 94.7%',
        },
        {
            'layout': 1,
            'title': 'Market Landscape',
            'body': 'Total addressable market expanded to $8.2B\nTop 3 competitors hold 45% combined share\nEmerging markets grew 32% in the quarter\nRegulatory changes favor our positioning',
        },
        {
            'layout': 1,
            'title': 'Product Portfolio Performance',
            'body': 'Enterprise Suite: $67M (+22% YoY)\nCloud Platform: $41M (+35% YoY)\nProfessional Services: $24M (+8% YoY)\nLegacy Products: $10M (-15% YoY)',
        },
        {
            'layout': 1,
            'title': 'Customer Acquisition Metrics',
            'body': 'New enterprise accounts: 47 (target: 40)\nAverage deal size increased to $285K\nSales cycle reduced from 92 to 78 days\nPartner-sourced deals: 31% of pipeline',
        },
        {
            'layout': 1,
            'title': 'Technology Roadmap',
            'body': 'AI-powered analytics module: Beta launch Oct 2025\nMobile platform redesign: GA Nov 2025\nAPI v3.0 with GraphQL support: Dec 2025\nSOC 2 Type II certification: Q4 2025',
        },
        {
            'layout': 1,
            'title': 'Financial Outlook',
            'body': 'Q4 revenue forecast: $155M-$162M\nFull-year guidance raised to $560M-$575M\nR&D investment increasing to 18% of revenue\nTarget operating margin: 25% by FY2026',
        },
        {
            'layout': 1,
            'title': 'Risk Factors',
            'body': 'Currency fluctuation impact: -$3.2M potential\nKey talent retention in competitive market\nSupply chain dependencies on 2 critical vendors\nPending regulatory review in EU markets',
        },
        {
            'layout': 1,
            'title': 'Strategic Initiatives',
            'body': 'Expand APAC sales team by 25 headcount\nLaunch vertical-specific solutions for healthcare\nEstablish innovation lab in Austin, TX\nPursue 2-3 strategic acquisitions under $50M',
        },
        {
            'layout': 1,
            'title': 'Next Steps & Action Items',
            'body': 'Board presentation scheduled for October 15\nBudget finalization due September 30\nPartner summit planning underway for November\nAnnual customer conference: January 2026',
        },
    ]

    for i, content in enumerate(slide_contents):
        layout_idx = content['layout']
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = content['title']

        # Set body/subtitle
        if layout_idx == 0 and 'subtitle' in content:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    ph.text = content['subtitle']
                    break
        elif 'body' in content:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    tf = ph.text_frame
                    tf.clear()
                    lines = content['body'].split('\n')
                    for j, line in enumerate(lines):
                        if j == 0:
                            tf.paragraphs[0].text = line
                        else:
                            p = tf.add_paragraph()
                            p.text = line
                    break

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Verify the master title placeholder position
    verify_prs = Presentation(OUTPUT)
    master = verify_prs.slide_masters[0]
    for ph in master.placeholders:
        if ph.placeholder_format.idx == 0:
            print(f'Master title placeholder: left={ph.left}, top={ph.top}, width={ph.width}, height={ph.height}')
            print(f'  In inches: left={ph.left/914400:.2f}, top={ph.top/914400:.2f}, width={ph.width/914400:.2f}')

    # Launch LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
