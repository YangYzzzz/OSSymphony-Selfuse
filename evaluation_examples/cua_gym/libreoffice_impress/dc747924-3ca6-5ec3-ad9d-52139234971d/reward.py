"""
Reward Script: Reset master slide title placeholder to centered top position
Task ID: impress_fix_050
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Master title placeholder left ~1.0in (centered horizontally)
  Component 2 (0.30): Master title placeholder top ~1.0in (at top of slide)
  Component 3 (0.25): Master title placeholder width ~8.0in
  Component 4 (0.15): Slide layouts inherit the corrected position
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_050'

# Expected values (from task context)
EXPECTED_LEFT = Inches(1)    # 914400 EMU
EXPECTED_TOP = Inches(1)     # 914400 EMU
EXPECTED_WIDTH = Inches(8)   # 7315200 EMU

# Tolerance: 5% relative or 0.1 inch absolute (whichever is larger)
def is_close(actual, expected, rel_tol=0.05, abs_tol=Inches(0.1)):
    """Check if actual EMU value is close to expected."""
    diff = abs(actual - expected)
    if diff <= abs_tol:
        return True
    if expected != 0 and diff / expected <= rel_tol:
        return True
    return False


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Get the slide master (should be only one)
    try:
        master = prs.slide_masters[0]
    except Exception as e:
        print(f"CRITICAL: Cannot access slide master: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Find the title placeholder on the master slide
    title_ph = None
    for ph in master.placeholders:
        if ph.placeholder_format.type == 1:  # TITLE type
            title_ph = ph
            break

    if title_ph is None:
        print("CRITICAL: No title placeholder found on master slide")
        print("REWARD: 0.0")
        return 0.0

    master_left = title_ph.left
    master_top = title_ph.top
    master_width = title_ph.width

    print(f"Master title placeholder: left={master_left} ({master_left/914400:.2f}in), "
          f"top={master_top} ({master_top/914400:.2f}in), "
          f"width={master_width} ({master_width/914400:.2f}in)")

    # Component 1: Master title placeholder left position ~1.0in (0.30 points)
    try:
        if is_close(master_left, EXPECTED_LEFT):
            print(f"PASS: Component 1 -- Master title left={master_left/914400:.2f}in, expected ~1.0in (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 1 -- Master title left={master_left/914400:.2f}in, expected ~1.0in")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Master title placeholder top position ~1.0in (0.30 points)
    try:
        if is_close(master_top, EXPECTED_TOP):
            print(f"PASS: Component 2 -- Master title top={master_top/914400:.2f}in, expected ~1.0in (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 2 -- Master title top={master_top/914400:.2f}in, expected ~1.0in")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Master title placeholder width ~8.0in (0.25 points)
    try:
        if is_close(master_width, EXPECTED_WIDTH):
            print(f"PASS: Component 3 -- Master title width={master_width/914400:.2f}in, expected ~8.0in (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 -- Master title width={master_width/914400:.2f}in, expected ~8.0in")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Slide layouts inherit the corrected position (0.15 points)
    # Check that at least some layouts have the corrected title placeholder position
    try:
        layouts_checked = 0
        layouts_correct = 0
        for layout in prs.slide_layouts:
            for ph in layout.placeholders:
                if ph.placeholder_format.type == 1:  # TITLE
                    layouts_checked += 1
                    if (is_close(ph.left, EXPECTED_LEFT) and
                        is_close(ph.top, EXPECTED_TOP) and
                        is_close(ph.width, EXPECTED_WIDTH)):
                        layouts_correct += 1
                    break  # only check first title PH per layout

        if layouts_checked > 0 and layouts_correct >= layouts_checked * 0.8:
            print(f"PASS: Component 4 -- {layouts_correct}/{layouts_checked} layouts have corrected title position (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 -- Only {layouts_correct}/{layouts_checked} layouts have corrected title position")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
