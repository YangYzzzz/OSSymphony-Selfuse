"""
FINAL REWARD SCRIPT - SUCCESS
Task: I’m tweaking slide 53 in LibreOffice Impress, and the first content text box looks a bit squeezed. How do I set its internal padding to exactly 0.3 cm on every side?
Generated: 2025-09-10 12:34:42
Status: success
Model: azure-o3
Total Steps: 5
"""

import os
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as PHT

# Constants
FILE_PATH = "/home/user/im_tweaking_slide_53_in_libreoffice_impress_and_the_first_content_text_box_looks_a_bit_squeezed_how__golden.pptx"
EMU_PER_INCH = 914400  # English Metric Units per inch
CM_PER_INCH = 2.54     # Centimetres per inch

# -----------------------------------------------------------------------------
# Utility helpers
# -----------------------------------------------------------------------------

def emu_to_cm(emu: int) -> float:
    """Convert EMU (English Metric Unit) to centimetres (float)."""
    return (emu / EMU_PER_INCH) * CM_PER_INCH


def verify_padding(measured_cm: float, target_cm: float = 0.3, tolerance_cm: float = 0.02) -> bool:
    """True if *measured_cm* is within ±*tolerance_cm* of *target_cm*."""
    return abs(measured_cm - target_cm) <= tolerance_cm


def find_first_content_shape(slide):
    """Return the first text–frame shape on *slide* that is NOT a title/subtitle."""
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        # Skip title-related placeholders
        placeholder_type = None
        if getattr(shape, "is_placeholder", False):
            placeholder_type = shape.placeholder_format.type
        if placeholder_type in {PHT.TITLE, PHT.CENTER_TITLE, PHT.SUBTITLE}:
            continue
        return shape  # Candidate content box found
    return None

# -----------------------------------------------------------------------------
# Main verification logic
# -----------------------------------------------------------------------------

def verify_task(file_path: str = FILE_PATH) -> float:
    """Verify that the first content text box on slide 53 has 0.3 cm padding on all sides."""
    total_score = 0.0
    max_score = 1.0
    print(f"Verifying internal padding for presentation: {file_path}")

    # --- 0. Prerequisites: file exists & loads (no points awarded) ---
    if not os.path.exists(file_path):
        print("✗ File not found – cannot proceed")
        return 0.0
    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded presentation containing {len(prs.slides)} slides")
    except Exception as exc:
        print(f"✗ Failed to load presentation: {exc}")
        return 0.0

    # --- 1. Ensure slide 53 exists ---
    if len(prs.slides) < 53:
        print("✗ Presentation has fewer than 53 slides – target slide missing")
        return 0.0
    print("✓ Slide 53 located")
    slide = prs.slides[52]  # zero-based index

    # --- 2. Locate first *content* text box (non-title) ---
    content_shape = find_first_content_shape(slide)
    if content_shape is None:
        print("✗ No suitable content text box found on slide 53")
        return 0.0
    print(f"✓ Found content text box: '{content_shape.name}'")

    # --- 3. Check each internal margin ---
    tf = content_shape.text_frame
    margins_emu = {
        "left": tf.margin_left,
        "right": tf.margin_right,
        "top": tf.margin_top,
        "bottom": tf.margin_bottom,
    }
    margins_cm = {side: emu_to_cm(val) for side, val in margins_emu.items()}

    # Report measured values
    for side, cm_val in margins_cm.items():
        print(f"  {side.capitalize()} margin: {cm_val:.3f} cm (target 0.3 cm)")

    # Award 0.25 for each correctly-set side
    for side, cm_val in margins_cm.items():
        if verify_padding(cm_val):
            print(f"    ✓ {side.capitalize()} margin within tolerance")
            total_score += 0.25
        else:
            print(f"    ✗ {side.capitalize()} margin outside tolerance")

    final_score = min(total_score, max_score)
    print(f"Total score: {final_score} (out of {max_score})")
    print(f"REWARD: {final_score}")
    return final_score

# -----------------------------------------------------------------------------
# Execute verification when run as a script
# -----------------------------------------------------------------------------
if __name__ == "__main__":
    verify_task()

