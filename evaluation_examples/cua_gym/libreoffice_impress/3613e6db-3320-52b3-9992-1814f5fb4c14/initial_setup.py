"""
Initial Setup: Create a 10-slide lobby display presentation with no auto-advance or kiosk settings.
Task ID: impress_fix_033
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_033'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tf


def add_multi_text(slide, left, top, width, height, lines, font_size=16, color=None):
    """Add a text box with multiple lines."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.space_after = Pt(4)
        if p.runs:
            p.runs[0].font.size = Pt(font_size)
            if color:
                p.runs[0].font.color.rgb = color
    return tf


def create_initial():
    prs = Presentation()
    # Standard 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Welcome / Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x00, 0x2B, 0x5C)
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(2),
                 "Welcome to Meridian Tower", font_size=44, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2), Inches(4), Inches(9), Inches(1.5),
                 "Your Premier Business & Innovation Hub", font_size=28,
                 color=RGBColor(0xCC, 0xD5, 0xE0), alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3), Inches(5.8), Inches(7), Inches(1),
                 "1200 Commerce Boulevard  |  San Francisco, CA 94105", font_size=16,
                 color=RGBColor(0x99, 0xAA, 0xBB), alignment=PP_ALIGN.CENTER)

    # --- Slide 2: Building Directory ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF0, 0xF2, 0xF5)
    add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
                 "Building Directory", font_size=36, bold=True,
                 color=RGBColor(0x00, 0x2B, 0x5C), alignment=PP_ALIGN.CENTER)
    directory_lines = [
        "Floor 1    |    Meridian Cafe  &  Visitor Center",
        "Floor 2    |    Apex Financial Advisors",
        "Floor 3    |    Blueline Legal Partners",
        "Floor 4    |    Stratos Marketing Group",
        "Floor 5    |    NovaTech Solutions",
        "Floor 6    |    Greenfield Architecture",
        "Floor 7    |    Pacific Health Associates",
        "Floor 8    |    SkyBridge Consulting",
        "Floor 9    |    Executive Suites & Conference Center",
        "Floor 10   |    Meridian Ventures (Rooftop Lounge)",
    ]
    add_multi_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(5),
                   directory_lines, font_size=18, color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 3: Today's Events ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
                 "Today's Events", font_size=36, bold=True,
                 color=RGBColor(0x00, 0x2B, 0x5C), alignment=PP_ALIGN.CENTER)
    events = [
        "9:00 AM   —   Board Meeting (Floor 9, Cascade Room)",
        "10:30 AM  —   Investor Pitch: NovaTech Q2 Preview (Floor 5)",
        "12:00 PM  —   Networking Lunch (Meridian Cafe, Floor 1)",
        "2:00 PM   —   Workshop: AI in Financial Services (Floor 2)",
        "3:30 PM   —   Design Review: Harbor District Project (Floor 6)",
        "5:00 PM   —   Happy Hour (Rooftop Lounge, Floor 10)",
    ]
    add_multi_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(4.5),
                   events, font_size=18, color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 4: Building Amenities ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x00, 0x2B, 0x5C)
    add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
                 "Building Amenities", font_size=36, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    amenities = [
        "Meridian Cafe — Artisan coffee & fresh pastries (Floor 1)",
        "Fitness Center — Open 5 AM – 10 PM, keycard access (Basement)",
        "Rooftop Lounge — Panoramic city views, open to all tenants",
        "Electric Vehicle Charging — 24 stations in parking garage",
        "Bike Storage & Showers — Secure racks for 60 bicycles",
        "Package Room — 24/7 secure parcel pick-up (Floor 1)",
    ]
    add_multi_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(4.5),
                   amenities, font_size=18, color=RGBColor(0xDD, 0xDD, 0xDD))

    # --- Slide 5: Visitor Information ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF0, 0xF2, 0xF5)
    add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
                 "Visitor Information", font_size=36, bold=True,
                 color=RGBColor(0x00, 0x2B, 0x5C), alignment=PP_ALIGN.CENTER)
    info = [
        "Check-in: All visitors must present ID at the front desk",
        "Wi-Fi: Connect to 'Meridian-Guest' (password: Welcome2025!)",
        "Parking: Visitor spots in Levels P1-P2 (first 2 hours free)",
        "Restrooms: Located near elevators on every floor",
        "Emergency: Dial 555 from any lobby phone",
        "Lost & Found: Contact front desk or call (415) 555-0199",
    ]
    add_multi_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(4.5),
                   info, font_size=18, color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 6: Sustainability Initiatives ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x5E, 0x20)
    add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
                 "Our Green Commitment", font_size=36, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    green = [
        "LEED Platinum Certified since 2021",
        "Solar panels generate 40% of building energy needs",
        "100% recycled water used for landscaping",
        "Zero-waste cafeteria program — 92% diversion rate",
        "Smart HVAC reduces energy consumption by 35%",
        "Annual carbon offset: 1,200 metric tons",
    ]
    add_multi_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(4.5),
                   green, font_size=18, color=RGBColor(0xE8, 0xF5, 0xE9))

    # --- Slide 7: Emergency Procedures ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xB7, 0x1C, 0x1C)
    add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
                 "Emergency Procedures", font_size=36, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    emergency = [
        "Fire: Evacuate via stairwells — DO NOT use elevators",
        "Assembly Point: North plaza, past the fountain",
        "Earthquake: Drop, Cover, Hold — move away from windows",
        "Medical: Call front desk (555) or 911",
        "Active Threat: Run, Hide, Fight — lock doors if sheltering",
        "AED Locations: Lobby, Floor 5, Floor 9 conference area",
    ]
    add_multi_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(4.5),
                   emergency, font_size=18, color=RGBColor(0xFF, 0xCC, 0xBC))

    # --- Slide 8: Upcoming Building Maintenance ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
                 "Upcoming Maintenance", font_size=36, bold=True,
                 color=RGBColor(0x00, 0x2B, 0x5C), alignment=PP_ALIGN.CENTER)
    maint = [
        "Apr 5 — Elevator B service (Floors 1-5), 10 PM – 6 AM",
        "Apr 8 — Fire alarm testing, brief alarms 9 AM – 11 AM",
        "Apr 12 — Window cleaning, Floors 7-10 (exterior only)",
        "Apr 15 — Parking garage restriping, Level P3 closed",
        "Apr 20 — HVAC filter replacement, minimal disruption",
        "Apr 28 — Annual generator test, 30-second power flicker at 3 AM",
    ]
    add_multi_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(4.5),
                   maint, font_size=18, color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 9: Community Board ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF0, 0xF2, 0xF5)
    add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
                 "Community Board", font_size=36, bold=True,
                 color=RGBColor(0x00, 0x2B, 0x5C), alignment=PP_ALIGN.CENTER)
    community = [
        "Yoga on the Rooftop — Every Wednesday 7 AM (free for tenants)",
        "Book Club — 'The Innovators' by Walter Isaacson, Apr 18 @ noon",
        "Charity Run: Bay to Breakers Team — Sign up at front desk",
        "Photography Contest: 'City at Dawn' — Submit by Apr 30",
        "Cooking Demo: Chef Maria Alvarez, Apr 22, Floor 1 Cafe",
        "Rideshare Board — Post commute offers at meridian-tower.com/rides",
    ]
    add_multi_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(4.5),
                   community, font_size=18, color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 10: Contact & Social ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x00, 0x2B, 0x5C)
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                 "Stay Connected", font_size=44, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    contact = [
        "Front Desk: (415) 555-0199",
        "Email: info@meridian-tower.com",
        "Web: www.meridian-tower.com",
        "Instagram: @MeridianTowerSF",
        "LinkedIn: Meridian Tower Management",
    ]
    add_multi_text(slide, Inches(2.5), Inches(3.5), Inches(8), Inches(3.5),
                   contact, font_size=20, color=RGBColor(0xCC, 0xD5, 0xE0))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
