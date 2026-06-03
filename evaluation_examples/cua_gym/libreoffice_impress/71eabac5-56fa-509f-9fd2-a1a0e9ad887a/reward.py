"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 153 needs a quick makeover: drop in a 6-row × 2-column comparison table, then shade rows 1, 3, and 5 with #E6E6E6 (10 % gray) while keeping rows 2, 4, and 6 white.
Generated: 2025-09-10 16:17:08
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE

FILE_PATH = "/home/user/slide_153_needs_a_quick_makeover_drop_in_a_6_row_2_column_comparison_table_then_shade_rows_1_3_and_5_golden.pptx"

def rgb_to_hex(rgb):
    """Convert python-pptx RGBColor to hex string like 'E6E6E6'."""
    return str(rgb).upper() if rgb else None

def locate_table(presentation, slide_idx):
    """Locate a 6×2 table on the given slide (0-based index)."""
    slide = presentation.slides[slide_idx]
    for shape in slide.shapes:
        if shape.has_table:
            tbl = shape.table
            rows, cols = len(tbl.rows), len(tbl.columns)
            if rows == 6 and cols == 2:
                print("✓ Found 6×2 table on slide 153")
                return tbl
            else:
                print(f"Table with wrong dimensions detected: {rows}×{cols}")
    print("✗ No 6×2 table found on slide 153")
    return None

def cell_fill_hex(cell):
    """Return hex string of cell fill or None if not RGB."""
    fc = cell.fill.fore_color
    if fc and fc.type == MSO_COLOR_TYPE.RGB and fc.rgb:
        return rgb_to_hex(fc.rgb)
    return None

def row_color_matches(cells, expected_hex):
    """Check that every cell in the row has the expected colour."""
    for cell in cells:
        seen = cell_fill_hex(cell)
        if expected_hex == "FFFFFF":  # white rows: accept explicit white or no fill
            if seen is None or seen == "FFFFFF":
                continue
            return False, seen
        else:  # grey rows must be explicit #E6E6E6
            if seen == expected_hex:
                continue
            return False, seen
    return True, None

def verify_presentation(path):
    print("Starting verification …")
    score = 0.0

    # 1. File existence & load (no points awarded – prerequisite)
    if not os.path.exists(path):
        print(f"✗ File not found: {path}")
        print("REWARD: 0.0")
        return 0.0
    try:
        prs = Presentation(path)
        print(f"✓ Loaded PPTX with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Could not load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # 2. Ensure slide 153 exists (no points – prerequisite)
    if len(prs.slides) < 153:
        print("✗ Presentation contains fewer than 153 slides")
        print("REWARD: 0.0")
        return 0.0

    # 3. Locate 6×2 table on slide 153 (0.30 points)
    tbl = locate_table(prs, 152)  # 0-based index
    if not tbl:
        print("Total score: 0.0")
        print("REWARD: 0.0")
        return 0.0
    score += 0.30

    # 4. Verify row shading
    odd_rows_ok = True
    even_rows_ok = True
    for r in range(6):
        expected_hex = "E6E6E6" if r % 2 == 0 else "FFFFFF"
        ok, found_hex = row_color_matches([tbl.cell(r, c) for c in range(2)], expected_hex)
        if ok:
            print(f"✓ Row {r + 1} colour correct")
        else:
            print(f"✗ Row {r + 1} colour incorrect – found {found_hex}, expected {expected_hex}")
            if r % 2 == 0:
                odd_rows_ok = False
            else:
                even_rows_ok = False

    # 4a. Odd rows correct? (0.35 points)
    if odd_rows_ok:
        score += 0.35
        print("✓ All odd rows shaded #E6E6E6 as required")
    else:
        print("✗ Some odd rows incorrectly shaded")

    # 4b. Even rows correct? (0.35 points)
    if even_rows_ok:
        score += 0.35
        print("✓ All even rows white or no-fill as required")
    else:
        print("✗ Some even rows incorrectly shaded")

    # Finalise score
    final_score = min(score, 1.0)
    # Account for floating-point precision
    if abs(final_score - 1.0) < 1e-6:
        final_score = 1.0

    print(f"Total score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score

if __name__ == "__main__":
    verify_presentation(FILE_PATH)

