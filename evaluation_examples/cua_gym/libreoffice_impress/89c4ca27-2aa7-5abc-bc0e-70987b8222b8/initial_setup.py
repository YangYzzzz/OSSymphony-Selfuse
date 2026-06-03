"""
Initial Setup: Design Portfolio presentation with 8 slides, dark gray backgrounds.
Slide 5 has title 'Project Spotlight' and empty content area.
Task ID: impress_ps_027
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_027'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

DARK_GRAY = RGBColor(0x2D, 0x2D, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB0, 0xB0)


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox


def create_initial():
    prs = Presentation()
    # Standard widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide1, DARK_GRAY)
    add_text_box(slide1, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
                 "Design Portfolio", font_size=44, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, Inches(1.5), Inches(3.8), Inches(10), Inches(1.0),
                 "Elena Vasquez  |  UI/UX Designer  |  2025", font_size=18,
                 color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # --- Slide 2: About Me ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide2, DARK_GRAY)
    add_text_box(slide2, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8),
                 "About Me", font_size=32, color=WHITE, bold=True)
    add_text_box(slide2, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.5),
                 "With over 8 years of experience in digital product design, "
                 "I specialize in creating intuitive user experiences for mobile and "
                 "web applications. My approach combines user research, interaction "
                 "design, and visual storytelling to deliver impactful solutions.",
                 font_size=14, color=LIGHT_GRAY)
    add_text_box(slide2, Inches(7.0), Inches(1.6), Inches(5.5), Inches(4.0),
                 "Skills: Figma, Sketch, Adobe XD, Prototyping, User Research, "
                 "Design Systems, Motion Design, Accessibility",
                 font_size=14, color=LIGHT_GRAY)

    # --- Slide 3: Project - E-Commerce Redesign ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide3, DARK_GRAY)
    add_text_box(slide3, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
                 "E-Commerce Platform Redesign", font_size=28, color=WHITE, bold=True)
    add_text_box(slide3, Inches(0.8), Inches(1.6), Inches(5.5), Inches(1.0),
                 "Client: RetailMax  |  2024", font_size=14, color=LIGHT_GRAY)
    add_text_box(slide3, Inches(0.8), Inches(2.8), Inches(11), Inches(3.5),
                 "Redesigned the entire shopping experience for RetailMax, resulting in "
                 "a 35% increase in conversion rates. The project involved restructuring "
                 "the product catalog, streamlining checkout flow, and implementing "
                 "a personalized recommendation engine.",
                 font_size=14, color=LIGHT_GRAY)

    # --- Slide 4: Project - Healthcare Dashboard ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide4, DARK_GRAY)
    add_text_box(slide4, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
                 "Healthcare Analytics Dashboard", font_size=28, color=WHITE, bold=True)
    add_text_box(slide4, Inches(0.8), Inches(1.6), Inches(5.5), Inches(1.0),
                 "Client: MedVista Health  |  2024", font_size=14, color=LIGHT_GRAY)
    add_text_box(slide4, Inches(0.8), Inches(2.8), Inches(11), Inches(3.5),
                 "Designed a comprehensive analytics dashboard for healthcare "
                 "professionals to monitor patient outcomes, track treatment efficacy, "
                 "and visualize population health trends in real time.",
                 font_size=14, color=LIGHT_GRAY)

    # --- Slide 5: Project Spotlight (TASK TARGET - empty content area) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide5, DARK_GRAY)
    add_text_box(slide5, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
                 "Project Spotlight", font_size=32, color=WHITE, bold=True)
    # Empty content area - no other content on this slide

    # --- Slide 6: Project - Travel App ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide6, DARK_GRAY)
    add_text_box(slide6, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
                 "Travel Companion App", font_size=28, color=WHITE, bold=True)
    add_text_box(slide6, Inches(0.8), Inches(1.6), Inches(5.5), Inches(1.0),
                 "Client: Wanderlust Inc  |  2023", font_size=14, color=LIGHT_GRAY)
    add_text_box(slide6, Inches(0.8), Inches(2.8), Inches(11), Inches(3.5),
                 "Created an immersive travel planning application with interactive "
                 "maps, AI-powered itinerary suggestions, and social sharing features "
                 "that garnered 500K+ downloads in the first quarter.",
                 font_size=14, color=LIGHT_GRAY)

    # --- Slide 7: Testimonials ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide7, DARK_GRAY)
    add_text_box(slide7, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
                 "Client Testimonials", font_size=32, color=WHITE, bold=True)
    add_text_box(slide7, Inches(0.8), Inches(1.8), Inches(5.0), Inches(2.0),
                 '"Elena transformed our vision into a product our users love. '
                 'Her attention to detail is unmatched."\n'
                 '- Sarah Kim, CEO of RetailMax',
                 font_size=14, color=LIGHT_GRAY)
    add_text_box(slide7, Inches(7.0), Inches(1.8), Inches(5.0), Inches(2.0),
                 '"Working with Elena was a seamless experience. She delivered '
                 'beyond expectations on every milestone."\n'
                 '- Dr. James Park, CTO of MedVista',
                 font_size=14, color=LIGHT_GRAY)

    # --- Slide 8: Contact ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide8, DARK_GRAY)
    add_text_box(slide8, Inches(1.5), Inches(2.0), Inches(10), Inches(1.0),
                 "Let's Work Together", font_size=36, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide8, Inches(1.5), Inches(3.5), Inches(10), Inches(2.5),
                 "elena.vasquez@designstudio.com\n"
                 "linkedin.com/in/elenavasquez\n"
                 "dribbble.com/elenavasquez\n"
                 "+1 (415) 555-0189",
                 font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
