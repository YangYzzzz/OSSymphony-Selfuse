"""
Reward Script: Project showcase slide verification for Design_Portfolio.pptx
Task ID: impress_ps_027
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25) - Rectangle placeholder on left side (~14x12cm, white border)
  Component 2 (0.25) - 'Mobile Banking App' title: 22pt bold white
  Component 3 (0.20) - 'FinTech Corp' and '2025' in 14pt light gray
  Component 4 (0.15) - Description paragraph: 12pt white, 3 lines
  Component 5 (0.15) - Right-side text positioned to the right of rectangle
"""

import os

from pptx import Presentation
from pptx.util import Pt, Emu

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_027'


def persist_app_state(domain: str):
    """Save any unsaved GUI edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # slide 5 (0-indexed)
    slide_width = prs.slide_width

    # Collect shapes on slide 5 (excluding the pre-existing 'Project Spotlight' title)
    shapes = list(slide.shapes)

    # Helper: find shapes by text content
    def find_text_shapes_containing(text_query):
        results = []
        for s in shapes:
            if s.has_text_frame:
                full_text = s.text_frame.text.strip()
                if text_query.lower() in full_text.lower():
                    results.append(s)
        return results

    # Helper: find rectangle/auto shapes (non-text-box)
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def find_rectangles():
        results = []
        for s in shapes:
            if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                results.append(s)
        return results

    # Component 1: Rectangle placeholder on left side with white border (0.25 points)
    try:
        rects = find_rectangles()
        rect_found = False
        if len(rects) > 0:
            for rect in rects:
                # Check approximate size: ~14cm wide (5040000 EMU), ~12cm tall (4320000 EMU)
                w_cm = rect.width / 360000
                h_cm = rect.height / 360000
                # Allow some tolerance: width 10-18cm, height 8-16cm
                size_ok = (10 <= w_cm <= 18) and (8 <= h_cm <= 16)
                # Must be on left side: left edge < 50% of slide width
                left_ok = rect.left < slide_width * 0.5

                # Check for white border/outline
                border_ok = False
                try:
                    line = rect.line
                    if line.fill.type is not None:
                        try:
                            rgb = str(line.color.rgb)
                            if rgb.upper() == "FFFFFF":
                                border_ok = True
                        except Exception:
                            pass
                except Exception:
                    pass

                if size_ok and left_ok and border_ok:
                    rect_found = True
                    print(f"PASS: Component 1 - Rectangle placeholder found: {w_cm:.1f}x{h_cm:.1f}cm, left side, white border (0.25 pts)")
                    total_score += 0.25
                    break

            if not rect_found:
                # Partial: rectangle exists but not all properties match
                for rect in rects:
                    w_cm = rect.width / 360000
                    h_cm = rect.height / 360000
                    left_ok = rect.left < slide_width * 0.5
                    print(f"FAIL: Component 1 - Rectangle found ({w_cm:.1f}x{h_cm:.1f}cm, left={left_ok}) but missing white border or wrong size/position")
                    break
        else:
            print(f"FAIL: Component 1 - No rectangle/auto shape found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: 'Mobile Banking App' in 22pt bold white (0.25 points)
    try:
        mba_shapes = find_text_shapes_containing("Mobile Banking App")
        if len(mba_shapes) > 0:
            found_correct = False
            for s in mba_shapes:
                for para in s.text_frame.paragraphs:
                    for run in para.runs:
                        if "Mobile Banking App" in run.text:
                            # Check font properties
                            is_bold = run.font.bold is True
                            size_ok = False
                            if run.font.size is not None:
                                size_pt = run.font.size / 12700  # EMU to pt
                                size_ok = abs(size_pt - 22) <= 2  # tolerance of 2pt
                            color_ok = False
                            try:
                                if run.font.color.type is not None:
                                    color_ok = str(run.font.color.rgb).upper() == "FFFFFF"
                            except Exception:
                                pass

                            if is_bold and size_ok and color_ok:
                                print(f"PASS: Component 2 - 'Mobile Banking App' found: bold, ~22pt, white (0.25 pts)")
                                total_score += 0.25
                                found_correct = True
                                break
                            else:
                                print(f"FAIL: Component 2 - 'Mobile Banking App' found but: bold={is_bold}, size_ok={size_ok}, color_ok={color_ok}")
                    if found_correct:
                        break
                if found_correct:
                    break
            if not found_correct and total_score < 0.25:
                pass  # already printed failure
        else:
            print(f"FAIL: Component 2 - 'Mobile Banking App' text not found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: 'FinTech Corp' and '2025' in ~14pt light gray (0.20 points)
    try:
        ft_found = False
        year_found = False

        # Check FinTech Corp
        ft_shapes = find_text_shapes_containing("FinTech Corp")
        for s in ft_shapes:
            for para in s.text_frame.paragraphs:
                for run in para.runs:
                    if "FinTech Corp" in run.text:
                        size_ok = False
                        if run.font.size is not None:
                            size_pt = run.font.size / 12700
                            size_ok = abs(size_pt - 14) <= 2
                        color_ok = False
                        try:
                            if run.font.color.type is not None:
                                rgb_val = str(run.font.color.rgb).upper()
                                # Light gray: B0B0B0 or similar gray shades
                                r, g, b = int(rgb_val[0:2], 16), int(rgb_val[2:4], 16), int(rgb_val[4:6], 16)
                                color_ok = (r > 140 and g > 140 and b > 140) and (r < 220 and g < 220 and b < 220)
                        except Exception:
                            pass
                        if size_ok and color_ok:
                            ft_found = True
                            break
                if ft_found:
                    break
            if ft_found:
                break

        # Check 2025
        year_shapes = find_text_shapes_containing("2025")
        for s in year_shapes:
            for para in s.text_frame.paragraphs:
                for run in para.runs:
                    if "2025" in run.text:
                        size_ok = False
                        if run.font.size is not None:
                            size_pt = run.font.size / 12700
                            size_ok = abs(size_pt - 14) <= 2
                        color_ok = False
                        try:
                            if run.font.color.type is not None:
                                rgb_val = str(run.font.color.rgb).upper()
                                r, g, b = int(rgb_val[0:2], 16), int(rgb_val[2:4], 16), int(rgb_val[4:6], 16)
                                color_ok = (r > 140 and g > 140 and b > 140) and (r < 220 and g < 220 and b < 220)
                        except Exception:
                            pass
                        if size_ok and color_ok:
                            year_found = True
                            break
                if year_found:
                    break
            if year_found:
                break

        if ft_found and year_found:
            print(f"PASS: Component 3 - 'FinTech Corp' and '2025' in ~14pt light gray (0.20 pts)")
            total_score += 0.20
        elif ft_found or year_found:
            print(f"PARTIAL: Component 3 - FinTech Corp={ft_found}, 2025={year_found} (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 3 - Neither 'FinTech Corp' nor '2025' found with correct formatting")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Description paragraph in ~12pt white, 3 lines (0.15 points)
    try:
        desc_found = False
        # Look for a text shape with longer text that mentions banking-related content
        for s in shapes:
            if not s.has_text_frame:
                continue
            full_text = s.text_frame.text.strip()
            # Description should be at least 50 chars and not be one of the other known texts
            if len(full_text) > 50 and "Mobile Banking App" not in full_text and "Project Spotlight" not in full_text:
                # Check that it has roughly 3 lines (line breaks or multi-line)
                # Count lines by line break chars or separate runs
                line_count = full_text.count('\x0b') + full_text.count('\n') + 1
                lines_ok = line_count >= 2  # at least 2 line breaks = 3 lines

                # Check font: first run should be ~12pt white
                first_run = None
                for para in s.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            first_run = run
                            break
                    if first_run:
                        break

                size_ok = False
                color_ok = False
                if first_run:
                    if first_run.font.size is not None:
                        size_pt = first_run.font.size / 12700
                        size_ok = abs(size_pt - 12) <= 2
                    try:
                        if first_run.font.color.type is not None:
                            color_ok = str(first_run.font.color.rgb).upper() == "FFFFFF"
                    except Exception:
                        pass

                if lines_ok and size_ok and color_ok:
                    print(f"PASS: Component 4 - Description found: {line_count} lines, ~12pt, white (0.15 pts)")
                    total_score += 0.15
                    desc_found = True
                    break
                elif lines_ok or (size_ok and color_ok):
                    print(f"PARTIAL: Component 4 - Description found but: lines_ok={lines_ok}, size_ok={size_ok}, color_ok={color_ok}")
                    # Small partial for having the description at all
                    break
                else:
                    print(f"FAIL: Component 4 - Long text found but formatting mismatch: lines_ok={lines_ok}, size_ok={size_ok}, color_ok={color_ok}")
                    break

        if not desc_found and total_score < 0.85:
            # Check if we even found any long text
            long_texts = [s for s in shapes if s.has_text_frame and len(s.text_frame.text.strip()) > 50
                          and "Project Spotlight" not in s.text_frame.text]
            if not long_texts:
                print(f"FAIL: Component 4 - No description paragraph found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Right-side text positioned to the right of the rectangle (0.15 points)
    try:
        # Find right-side text shapes: 'Mobile Banking App', 'FinTech Corp', '2025', description
        # They should be positioned with left > 50% of slide width
        right_text_keywords = ["Mobile Banking App", "FinTech Corp", "2025"]
        right_count = 0
        total_expected = 3  # 3 text items we check

        for keyword in right_text_keywords:
            kw_shapes = find_text_shapes_containing(keyword)
            for s in kw_shapes:
                if s.left > slide_width * 0.4:  # right side (with tolerance)
                    right_count += 1
                    break

        # Also check the description text
        for s in shapes:
            if s.has_text_frame and len(s.text_frame.text.strip()) > 50 and "Project Spotlight" not in s.text_frame.text:
                if s.left > slide_width * 0.4:
                    right_count += 1
                break
        total_expected = 4

        if right_count >= 3:
            print(f"PASS: Component 5 - {right_count}/{total_expected} text elements on right side (0.15 pts)")
            total_score += 0.15
        elif right_count >= 1:
            partial = round(0.15 * right_count / total_expected, 2)
            print(f"PARTIAL: Component 5 - Only {right_count}/{total_expected} text elements on right side ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 - No text elements positioned on right side of slide")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
