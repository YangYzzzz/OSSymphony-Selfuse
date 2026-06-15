"""
Reward Script: Add footer to slides 2-10, hide on title slide
Task ID: impress_teach_011
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Footer presence - at least 5 of slides 2-10 have a footer placeholder
  Component 2 (0.3): Footer text exact - all 9 slides (2-10) have exact footer text
  Component 3 (0.3): Title slide exclusion - slide 1 has no footer AND at least one other slide does
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_011'
EXPECTED_FOOTER = 'CS101 - Introduction to Computer Science - Confidential'


def persist_app_state(domain: str):
    """Save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_footer_text(slide):
    """
    Extract footer text from a slide by finding a placeholder with idx=11 (FOOTER type)
    or any placeholder whose name contains 'Footer'.
    Returns the footer text string, or None if no footer found.
    """
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            # Footer placeholder has idx=11 in standard PPTX
            if shape.placeholder_format.idx == 11:
                if shape.has_text_frame:
                    return shape.text_frame.text.strip()
                return ""
        # Also check by name as fallback
        if hasattr(shape, 'name') and 'footer' in (shape.name or '').lower():
            if shape.has_text_frame:
                return shape.text_frame.text.strip()
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    if num_slides < 2:
        print(f"FAIL: Expected at least 10 slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    # Collect footer info for all slides
    footer_info = {}
    for i, slide in enumerate(prs.slides):
        footer_text = get_footer_text(slide)
        footer_info[i + 1] = footer_text  # 1-indexed slide number
        print(f"  Slide {i+1}: footer = {repr(footer_text)}")

    # Component 1: Footer presence on slides 2-10 (0.4 points)
    # At least 5 of slides 2-10 must have a footer placeholder containing "Confidential"
    try:
        slides_with_footer = 0
        for slide_num in range(2, min(num_slides + 1, 11)):
            ft = footer_info.get(slide_num)
            if ft is not None and 'Confidential' in ft:
                slides_with_footer += 1

        if slides_with_footer >= 5:
            print(f"PASS: Component 1 — {slides_with_footer}/9 slides (2-10) have footer with 'Confidential' (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 — only {slides_with_footer}/9 slides (2-10) have footer with 'Confidential', need >= 5")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Footer text exact match on all slides 2-10 (0.3 points)
    # All 9 slides must have the exact expected footer text
    try:
        exact_match_count = 0
        for slide_num in range(2, min(num_slides + 1, 11)):
            ft = footer_info.get(slide_num)
            if ft == EXPECTED_FOOTER:
                exact_match_count += 1

        expected_count = min(num_slides, 10) - 1  # slides 2 through min(N, 10)
        if exact_match_count == expected_count:
            print(f"PASS: Component 2 — all {exact_match_count} slides (2-10) have exact footer text (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 — {exact_match_count}/{expected_count} slides have exact footer text '{EXPECTED_FOOTER}'")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Title slide exclusion (0.3 points)
    # Slide 1 must NOT have a footer AND at least one of slides 2-10 must have a footer
    # This compound check ensures it fails on initial_env (where no slides have footers)
    try:
        slide1_footer = footer_info.get(1)
        slide1_has_no_footer = (slide1_footer is None)
        any_other_has_footer = any(
            footer_info.get(s) is not None
            for s in range(2, min(num_slides + 1, 11))
        )

        if slide1_has_no_footer and any_other_has_footer:
            print(f"PASS: Component 3 — slide 1 has no footer, other slides have footers (0.3 pts)")
            total_score += 0.3
        elif not slide1_has_no_footer:
            print(f"FAIL: Component 3 — slide 1 has footer: {repr(slide1_footer)}")
        else:
            print(f"FAIL: Component 3 — no other slides have footer (title exclusion requires footers on other slides)")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
