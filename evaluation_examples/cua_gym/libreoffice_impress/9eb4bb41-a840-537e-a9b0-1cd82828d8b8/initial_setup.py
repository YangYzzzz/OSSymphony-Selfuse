"""
Initial Setup: CS101 Lecture Presentation with 10 slides, no footers
Task ID: impress_teach_011
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_011'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=None, alignment=None):
    """Helper to add a text box with formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()

    # Slide dimensions (standard 16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "CS101 - Introduction to Computer Science"
    slide1.placeholders[1].text = "Fall 2025 Semester\nProfessor Elena Rodriguez\nDepartment of Computer Science"

    # --- Slide 2: Course Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Course Overview"
    tf = slide2.placeholders[1].text_frame
    tf.text = "Course Objectives"
    p = tf.add_paragraph()
    p.text = "Understand fundamental concepts of computer science"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Develop problem-solving skills using computational thinking"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Learn basic programming constructs in Python"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Explore data structures, algorithms, and software design"
    p.level = 1

    # --- Slide 3: What is Computer Science? ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "What is Computer Science?"
    tf = slide3.placeholders[1].text_frame
    tf.text = "Computer science is the study of computation, information, and automation."
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Key Areas:"
    p.runs[0].font.bold = True
    for item in ["Theoretical Computer Science", "Computer Systems & Architecture",
                  "Artificial Intelligence & Machine Learning", "Software Engineering",
                  "Data Science & Analytics"]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 4: History of Computing ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "History of Computing"
    tf = slide4.placeholders[1].text_frame
    tf.text = "1830s - Charles Babbage designs the Analytical Engine"
    p = tf.add_paragraph()
    p.text = "1936 - Alan Turing introduces the Turing Machine concept"
    p = tf.add_paragraph()
    p.text = "1945 - ENIAC, the first general-purpose electronic computer"
    p = tf.add_paragraph()
    p.text = "1969 - ARPANET connects four universities"
    p = tf.add_paragraph()
    p.text = "1991 - Tim Berners-Lee launches the World Wide Web"
    p = tf.add_paragraph()
    p.text = "2007 - The smartphone revolution begins"

    # --- Slide 5: Binary and Data Representation ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Binary and Data Representation"
    tf = slide5.placeholders[1].text_frame
    tf.text = "All digital information is stored using binary (0s and 1s)"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Common Units:"
    p.runs[0].font.bold = True
    for item in ["1 Bit = 0 or 1", "1 Byte = 8 Bits",
                  "1 Kilobyte (KB) = 1,024 Bytes", "1 Megabyte (MB) = 1,024 KB",
                  "1 Gigabyte (GB) = 1,024 MB"]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 6: Introduction to Algorithms ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Introduction to Algorithms"
    tf = slide6.placeholders[1].text_frame
    tf.text = "An algorithm is a step-by-step procedure for solving a problem."
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Properties of a Good Algorithm:"
    p.runs[0].font.bold = True
    for item in ["Correctness - produces the right output",
                  "Efficiency - uses minimal resources",
                  "Clarity - easy to understand and implement",
                  "Finiteness - terminates after a finite number of steps"]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 7: Programming Fundamentals ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Programming Fundamentals"
    tf = slide7.placeholders[1].text_frame
    tf.text = "Core Programming Concepts"
    p = tf.add_paragraph()
    p.text = "Variables and Data Types (int, float, str, bool)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Control Flow (if/else, loops, functions)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Input/Output Operations"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Error Handling and Debugging"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Code Organization and Modularity"
    p.level = 1

    # --- Slide 8: Data Structures ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Data Structures"
    tf = slide8.placeholders[1].text_frame
    tf.text = "Fundamental Data Structures:"
    p = tf.add_paragraph()
    p.text = "Arrays - fixed-size, indexed collections"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Linked Lists - dynamic, node-based sequences"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Stacks - Last In, First Out (LIFO)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Queues - First In, First Out (FIFO)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Trees - hierarchical parent-child relationships"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Hash Tables - key-value pair storage"
    p.level = 1

    # --- Slide 9: Software Development Lifecycle ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Software Development Lifecycle"
    tf = slide9.placeholders[1].text_frame
    tf.text = "Phases of Software Development:"
    for item in ["1. Requirements Gathering and Analysis",
                  "2. System Design and Architecture",
                  "3. Implementation and Coding",
                  "4. Testing and Quality Assurance",
                  "5. Deployment and Release",
                  "6. Maintenance and Iteration"]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 10: Course Schedule & Assessment ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "Course Schedule & Assessment"
    tf = slide10.placeholders[1].text_frame
    tf.text = "Assessment Breakdown:"
    for item in ["Homework Assignments: 30%",
                  "Midterm Exam: 20%",
                  "Final Project: 25%",
                  "Final Exam: 20%",
                  "Class Participation: 5%"]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Office Hours: Tuesdays & Thursdays, 2:00 PM - 4:00 PM, Room 312"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
