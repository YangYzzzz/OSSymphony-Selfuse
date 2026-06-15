"""
Reward Script: Configure all slides with consistent textbox sizing (36pt/18pt)
Task ID: osworld_impress_textbox_fontsize_specific_009
Domain: libreoffice_impress
Scoring:
  - Component 1: First textbox on each slide is 36pt (0.1 per slide, 0.5 total)
  - Component 2: Second textbox on each slide is 18pt (0.1 per slide, 0.5 total)
"""

import os

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_textbox_fontsize_specific_009'

EXPECTED_FIRST_PT = 36.0
EXPECTED_SECOND_PT = 18.0
NUM_SLIDES = 5
PER_SLIDE_SCORE = 0.1  # 0.1 per slide per component


def get_run_font_size_pt(shape):
    """
    Return the font size (in points) for the first non-empty run
    found in the text frame of a shape.
    EMU to Pt conversion: 1 pt = 12700 EMU.
    Returns None if no run has an explicit size.
    """
    if not shape.has_text_frame:
        return None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                return run.font.size / 12700.0
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation — abort early on failure
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: verify we have exactly 5 slides
    if len(prs.slides) != NUM_SLIDES:
        print(f"PRECONDITION FAIL: Expected {NUM_SLIDES} slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: First textbox on each slide is 36pt (0.1 per slide, 0.5 total)
    # Correctly fails on initial_env (first textbox sizes: 28, 32, 24, 30, 26 pt)
    # Correctly passes on golden_env (all first textboxes = 36 pt)
    comp1_pass = 0
    for slide_idx, slide in enumerate(prs.slides):
        try:
            text_shapes = [s for s in slide.shapes if s.has_text_frame]
            if len(text_shapes) < 1:
                print(f"FAIL: Component 1 — Slide {slide_idx + 1} has no textboxes")
                continue

            first_shape = text_shapes[0]
            actual_pt = get_run_font_size_pt(first_shape)

            if actual_pt is not None and abs(actual_pt - EXPECTED_FIRST_PT) < 0.1:
                print(f"PASS: Slide {slide_idx + 1} first textbox = {actual_pt}pt (expected {EXPECTED_FIRST_PT}pt)")
                comp1_pass += 1
                total_score += PER_SLIDE_SCORE
            else:
                print(f"FAIL: Slide {slide_idx + 1} first textbox = {actual_pt}pt (expected {EXPECTED_FIRST_PT}pt)")
        except Exception as e:
            print(f"ERROR: Component 1 — Slide {slide_idx + 1}: {e}")

    print(f"Component 1 (first textbox 36pt): {comp1_pass}/{NUM_SLIDES} slides correct — {comp1_pass * PER_SLIDE_SCORE:.1f} pts")

    # Component 2: Second textbox on each slide is 18pt (0.1 per slide, 0.5 total)
    # Correctly fails on initial_env (second textbox sizes: 22, 14, 20, 16, 24 pt)
    # Correctly passes on golden_env (all second textboxes = 18 pt)
    comp2_pass = 0
    for slide_idx, slide in enumerate(prs.slides):
        try:
            text_shapes = [s for s in slide.shapes if s.has_text_frame]
            if len(text_shapes) < 2:
                print(f"FAIL: Component 2 — Slide {slide_idx + 1} has fewer than 2 textboxes")
                continue

            second_shape = text_shapes[1]
            actual_pt = get_run_font_size_pt(second_shape)

            if actual_pt is not None and abs(actual_pt - EXPECTED_SECOND_PT) < 0.1:
                print(f"PASS: Slide {slide_idx + 1} second textbox = {actual_pt}pt (expected {EXPECTED_SECOND_PT}pt)")
                comp2_pass += 1
                total_score += PER_SLIDE_SCORE
            else:
                print(f"FAIL: Slide {slide_idx + 1} second textbox = {actual_pt}pt (expected {EXPECTED_SECOND_PT}pt)")
        except Exception as e:
            print(f"ERROR: Component 2 — Slide {slide_idx + 1}: {e}")

    print(f"Component 2 (second textbox 18pt): {comp2_pass}/{NUM_SLIDES} slides correct — {comp2_pass * PER_SLIDE_SCORE:.1f} pts")

    final_score = round(min(total_score, 1.0), 1)
    print(f"\nScore: {total_score:.1f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in a given env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
