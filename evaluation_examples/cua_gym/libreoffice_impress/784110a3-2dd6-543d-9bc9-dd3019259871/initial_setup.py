"""
Initial Setup: Configure all slides in this deck with consistent textbox sizing
Task ID: osworld_impress_textbox_fontsize_specific_009
Domain: libreoffice_impress

Creates a 5-slide presentation where each slide has two textboxes (title + body)
with INCONSISTENT font sizes. The agent must make them consistent (36pt / 18pt).
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_textbox_fontsize_specific_009'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Each slide: inconsistent font sizes - never (36pt title, 18pt body) simultaneously
    # Slide data: (title_text, body_text, title_font_pt, body_font_pt)
    slides_data = [
        (
            "Q1 Business Review",
            "Revenue grew by 12% year-over-year, driven by strong performance in the North American market and new enterprise contract wins across the EMEA region.",
            28,  # title: should be 36
            22,  # body: should be 18
        ),
        (
            "Product Roadmap 2025",
            "Key milestones include the launch of our mobile platform in March, the API v3 release in June, and the enterprise dashboard rollout scheduled for Q4.",
            32,  # title: should be 36
            14,  # body: should be 18
        ),
        (
            "Team Achievements",
            "Our engineering team delivered 47 feature releases with a 98.7% uptime record. Customer satisfaction scores reached an all-time high of 4.8 out of 5.",
            24,  # title: should be 36
            20,  # body: should be 18
        ),
        (
            "Market Expansion Strategy",
            "Targeting Southeast Asia and Latin America as primary growth markets for 2025. Initial investment of $2.4M allocated for regional office setup and local hiring.",
            30,  # title: should be 36
            16,  # body: should be 18
        ),
        (
            "Next Steps & Action Items",
            "All department heads to submit updated budget proposals by February 15. Quarterly review meetings to be scheduled with board members no later than March 1.",
            26,  # title: should be 36
            24,  # body: should be 18
        ),
    ]

    blank_layout = prs.slide_layouts[6]  # Blank layout

    for title_text, body_text, title_pt, body_pt in slides_data:
        slide = prs.slides.add_slide(blank_layout)

        # --- First textbox: Title ---
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9.0), Inches(1.2)
        )
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.alignment = PP_ALIGN.LEFT
        run_title = p_title.add_run()
        run_title.text = title_text
        run_title.font.size = Pt(title_pt)
        run_title.font.bold = True
        run_title.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)

        # --- Second textbox: Body ---
        body_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.8), Inches(9.0), Inches(4.5)
        )
        tf_body = body_box.text_frame
        tf_body.word_wrap = True
        p_body = tf_body.paragraphs[0]
        p_body.alignment = PP_ALIGN.LEFT
        run_body = p_body.add_run()
        run_body.text = body_text
        run_body.font.size = Pt(body_pt)
        run_body.font.bold = False
        run_body.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
