"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 32—the wrap-up slide—I need to swap the colors so it pops on the projector: make the entire slide background solid Dark Blue 1 (#002060) and change the title text itself to pure white (#FFFFFF) in LibreOffice Impress. What steps get that done?
Generated: 2025-09-10 22:45:10
Status: success
Model: azure-o3
Total Steps: 1
"""

from pptx import Presentation
from pptx.dml.color import RGBColor, MSO_THEME_COLOR
import os, traceback

"""
Reward Script: Verify color swap on slide 32
Task Requirements:
1. Slide 32 background must be solid Dark Blue 1 (#002060)
2. Title text on slide 32 must be pure white (#FFFFFF)

Scoring:
- 0.5 points for correct background color
- 0.5 points for all title text runs being white
Returns progressive score from 0.0 to 1.0 and prints detailed verification steps.
"""

FILE_PATH = (
    "/home/user/on_slide_32the_wrap_up_slidei_need_to_swap_the_colors_so_it_pops_on_the_projector_"
    "make_the_entire_sl_golden.pptx"
)


def is_dark_blue1(color_format):
    """Check if ColorFormat equals #002060 (RGB 0,32,96)."""
    if color_format is None:
        return False
    try:
        if color_format.type == 1:  # RGB
            return color_format.rgb == RGBColor(0, 32, 96)
    except Exception:
        pass
    return False


def is_white(color_format):
    """Check if ColorFormat equals white (#FFFFFF) via RGB or theme-equivalent."""
    if color_format is None:
        return False
    try:
        if color_format.type == 1 and color_format.rgb == RGBColor(255, 255, 255):
            return True
        if (
            hasattr(color_format, "theme_color")
            and color_format.theme_color
            in {MSO_THEME_COLOR.BACKGROUND_1, MSO_THEME_COLOR.LIGHT_1}
        ):
            return True
    except Exception:
        pass
    return False


def verify_slide_colors(file_path):
    print(f"Loading presentation: {file_path}")
    if not os.path.exists(file_path):
        print("✗ File not found.")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    target_idx = 31  # zero-based index for slide 32
    if len(prs.slides) <= target_idx:
        print(
            f"✗ Not enough slides. Found {len(prs.slides)}, expected at least 32."
        )
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[target_idx]
    score = 0.0

    # 1. Verify background color
    try:
        fill = slide.background.fill
        bg_ok = False
        if fill.type == 1:  # Solid fill
            bg_ok = is_dark_blue1(fill.fore_color)
        else:
            if hasattr(fill, "fore_color"):
                bg_ok = is_dark_blue1(fill.fore_color)
        if bg_ok:
            print("✓ Background color is Dark Blue 1 (#002060). (+0.5)")
            score += 0.5
        else:
            print("✗ Background color is not Dark Blue 1 (#002060).")
    except Exception as e:
        print("✗ Error checking background color:", e)
        traceback.print_exc()

    # 2. Verify title text color
    title_shapes = [
        s
        for s in slide.shapes
        if s.has_text_frame
        and (
            (hasattr(s, "is_placeholder") and s.is_placeholder and s.placeholder_format.type in {1, 3})
            or ("title" in (s.name or "").lower())
        )
    ]

    if not title_shapes:
        print("✗ No title shape detected on slide 32.")
    else:
        titles_white = True
        for shp in title_shapes:
            for para in shp.text_frame.paragraphs:
                for run in para.runs:
                    if not is_white(run.font.color):
                        titles_white = False
                        print(
                            f"  ✗ Non-white text found: '{run.text}' with color {run.font.color.rgb if run.font.color.type==1 else run.font.color.theme_color}"
                        )
        if titles_white:
            print("✓ All title text runs are white (#FFFFFF). (+0.5)")
            score += 0.5
        else:
            print("✗ Not all title text runs are white (#FFFFFF).")

    final = round(min(score, 1.0), 2)
    print(f"REWARD: {final}")
    return final


if __name__ == "__main__":
    verify_slide_colors(FILE_PATH)

