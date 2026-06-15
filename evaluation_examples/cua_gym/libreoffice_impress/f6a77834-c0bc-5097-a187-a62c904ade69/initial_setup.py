"""
Initial Setup: 16-slide navigation deck with table of contents
Task ID: impress_fix_093
Domain: libreoffice_impress

Creates a presentation with 16 slides. Slide 1 is a Table of Contents.
Slides 2-16 contain various department/topic content. No navigation links exist.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_093'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


# Slide topics for slides 2-16
SLIDE_TOPICS = [
    ("Company Overview", [
        "Founded in 2018 by Dr. Elena Vasquez and James Park",
        "Headquarters in Austin, TX with offices in London and Singapore",
        "Over 2,400 employees across 12 countries",
        "Annual revenue exceeded $340M in FY2025",
    ]),
    ("Q1 2026 Financial Summary", [
        "Revenue: $92.3M (up 18% YoY)",
        "Operating margin improved to 24.7%",
        "Free cash flow: $31.5M",
        "Customer acquisition cost reduced by 12%",
    ]),
    ("Product Roadmap", [
        "v3.0 Launch: AI-powered analytics dashboard (May 2026)",
        "Mobile app redesign targeting 40% engagement lift",
        "Enterprise SSO and SCIM provisioning",
        "Real-time collaboration features in Q3",
    ]),
    ("Engineering Department", [
        "Team size: 680 engineers across 8 squads",
        "Migrated 94% of services to Kubernetes",
        "Average deploy frequency: 47 deploys/day",
        "Incident response SLA met 99.2% of the time",
    ]),
    ("Marketing & Growth", [
        "Launched 'Elevate' brand campaign in 14 markets",
        "Website traffic grew 62% quarter-over-quarter",
        "Content marketing drove 3,200 qualified leads",
        "NPS score improved from 42 to 58",
    ]),
    ("Sales Performance", [
        "Closed 187 enterprise deals in Q1",
        "Average deal size increased to $124K",
        "Sales cycle reduced from 68 to 51 days",
        "Pipeline coverage ratio: 3.8x",
    ]),
    ("Customer Success Metrics", [
        "Net retention rate: 118%",
        "Average CSAT score: 4.6/5.0",
        "Onboarding time reduced to 14 days",
        "Support ticket resolution: 4.2 hours average",
    ]),
    ("Human Resources Update", [
        "Hired 312 new employees in Q1 2026",
        "Employee satisfaction survey: 87% favorable",
        "Launched mentorship program with 150 pairs",
        "Voluntary attrition rate: 8.3% (industry avg: 13%)",
    ]),
    ("IT Infrastructure", [
        "Cloud spend optimized: saved $1.2M annually",
        "Zero-trust security model fully deployed",
        "99.99% uptime across production systems",
        "Completed SOC 2 Type II recertification",
    ]),
    ("Research & Innovation", [
        "Filed 23 patent applications in Q1",
        "Published 8 peer-reviewed papers",
        "Partnered with MIT Media Lab on NLP research",
        "Internal hackathon yielded 3 product features",
    ]),
    ("Global Expansion", [
        "Opened new offices in Tokyo and Sao Paulo",
        "Localized product in 6 new languages",
        "APAC revenue grew 34% year-over-year",
        "Established partnerships with 12 regional resellers",
    ]),
    ("Sustainability Initiatives", [
        "Achieved carbon neutral operations in Q4 2025",
        "100% renewable energy for all data centers",
        "Reduced paper usage by 78% company-wide",
        "Planted 50,000 trees through reforestation program",
    ]),
    ("Risk Management", [
        "Completed annual enterprise risk assessment",
        "Cybersecurity insurance coverage increased to $50M",
        "Business continuity plan tested successfully",
        "Regulatory compliance score: 97/100",
    ]),
    ("Strategic Partnerships", [
        "Signed co-development agreement with Salesforce",
        "AWS Advanced Technology Partner status achieved",
        "Joint GTM initiative with Snowflake",
        "Integration marketplace grew to 240+ connectors",
    ]),
    ("Next Steps & Action Items", [
        "Finalize Q2 budget allocations by April 15",
        "Launch v3.0 beta program with 50 pilot customers",
        "Complete Series D fundraising preparation",
        "Board meeting scheduled for May 8, 2026",
    ]),
]


def create_initial():
    prs = Presentation()
    # Standard widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Table of Contents ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Table of Contents"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # Subtitle
    sub_box = slide1.shapes.add_textbox(Inches(1), Inches(1.6), Inches(11), Inches(0.6))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "NovaTech Solutions — Q1 2026 Business Review"
    p_sub.alignment = PP_ALIGN.CENTER
    run_sub = p_sub.runs[0]
    run_sub.font.name = "Calibri"
    run_sub.font.size = Pt(18)
    run_sub.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # TOC entries in two columns
    left_col_topics = SLIDE_TOPICS[:8]
    right_col_topics = SLIDE_TOPICS[8:]

    # Left column
    left_box = slide1.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(5), Inches(4.5))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    for i, (topic, _) in enumerate(left_col_topics):
        if i == 0:
            p = tf_left.paragraphs[0]
        else:
            p = tf_left.add_paragraph()
        p.text = f"{i + 2}.  {topic}"
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.name = "Calibri"
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Right column
    right_box = slide1.shapes.add_textbox(Inches(7), Inches(2.5), Inches(5), Inches(4.5))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    for i, (topic, _) in enumerate(right_col_topics):
        if i == 0:
            p = tf_right.paragraphs[0]
        else:
            p = tf_right.add_paragraph()
        p.text = f"{i + 10}.  {topic}"
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.name = "Calibri"
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slides 2-16: Content slides ---
    for idx, (topic_title, bullets) in enumerate(SLIDE_TOPICS):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

        # Slide number indicator
        num_box = slide.shapes.add_textbox(Inches(11.5), Inches(0.3), Inches(1.5), Inches(0.5))
        tf_num = num_box.text_frame
        p_num = tf_num.paragraphs[0]
        p_num.text = f"Slide {idx + 2}"
        p_num.alignment = PP_ALIGN.RIGHT
        r_num = p_num.runs[0]
        r_num.font.name = "Calibri"
        r_num.font.size = Pt(12)
        r_num.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

        # Title
        t_box = slide.shapes.add_textbox(Inches(1), Inches(0.6), Inches(11), Inches(1.0))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = topic_title
        p_t.alignment = PP_ALIGN.LEFT
        r_t = p_t.runs[0]
        r_t.font.name = "Calibri"
        r_t.font.size = Pt(32)
        r_t.font.bold = True
        r_t.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

        # Divider line (thin rectangle)
        from pptx.enum.shapes import MSO_SHAPE
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(1.7), Inches(11), Inches(0.03)
        )

        # Bullet content
        b_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.5), Inches(4.0))
        tf_b = b_box.text_frame
        tf_b.word_wrap = True
        for bi, bullet in enumerate(bullets):
            if bi == 0:
                p_b = tf_b.paragraphs[0]
            else:
                p_b = tf_b.add_paragraph()
            p_b.text = f"•  {bullet}"
            p_b.space_after = Pt(12)
            r_b = p_b.runs[0]
            r_b.font.name = "Calibri"
            r_b.font.size = Pt(20)
            r_b.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
