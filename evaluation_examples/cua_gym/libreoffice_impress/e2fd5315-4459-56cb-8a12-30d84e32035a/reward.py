"""
FINAL REWARD SCRIPT - SUCCESS
Task: Set Table 1 column widths to 3.0 cm, 2.0 cm, 4.0 cm for the first three columns.
Generated: 2025-10-17 08:24:58
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# -------------------------------------------------------------
# Reward Script : Set Table 1 column widths to 3.0 cm, 2.0 cm, 4.0 cm
# -------------------------------------------------------------
# This script verifies that, in the provided presentation file,
# the FIRST table found (Table 1) has its first three columns
# set to exactly 3.0 cm, 2.0 cm, and 4.0 cm (±0.1 cm tolerance).
# Progressive scoring is awarded per-column so that partial
# credit is possible if only some widths are correct.
# -------------------------------------------------------------

FILE_PATH = "/home/user/set_table_1_column_widths_to_30_cm_20_cm_40_cm_for_the_first_three_columns.pptx"

# Constants
CM_TO_EMU = 360000                # 1 cm  -> 360,000 EMU
EXPECTED_CM = [3.0, 2.0, 4.0]     # Target widths in cm
EXPECTED_EMU = [int(round(cm * CM_TO_EMU)) for cm in EXPECTED_CM]
TOLERANCE_EMU = int(0.1 * CM_TO_EMU)  # ±0.1 cm tolerance
COLUMN_SCORES = [0.34, 0.33, 0.33]    # Progressive scoring split (sums to 1.0)


def verify_table_column_widths(file_path: str) -> float:
    """Verify the first table's first three column widths.

    Returns a float between 0.0 and 1.0 representing the reward.
    """
    total_score = 0.0
    max_score = 1.0

    # 1) File existence & loading (no points awarded, prerequisite only)
    if not os.path.exists(file_path):
        print("✗ File not found:", file_path)
        print("REWARD: 0.0")
        return 0.0
    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded presentation with {len(prs.slides)} slides (0 points – prerequisite)")
    except Exception as e:
        print("✗ Error loading presentation:", e)
        print("REWARD: 0.0")
        return 0.0

    # 2) Locate the first table (Table 1)
    first_table = None
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                first_table = shape.table
                print(f"✓ Found Table 1 on slide {slide_idx}, shape {shape_idx}")
                break
        if first_table:
            break

    if not first_table:
        print("✗ No table found in the presentation – cannot verify column widths")
        print("REWARD: 0.0")
        return 0.0

    # 3) Verify at least three columns are present
    if len(first_table.columns) < 3:
        print(f"✗ Table has only {len(first_table.columns)} column(s); requires at least 3")
        print("REWARD: 0.0")
        return 0.0

    # 4) Measure widths and compare against expectations
    measured_widths = [first_table.columns[i].width for i in range(3)]
    measured_cm = [round(w / CM_TO_EMU, 3) for w in measured_widths]

    print("Measured widths (EMU):", measured_widths)
    print("Measured widths (cm): ", measured_cm)

    for idx in range(3):
        diff = abs(measured_widths[idx] - EXPECTED_EMU[idx])
        if diff <= TOLERANCE_EMU:
            total_score += COLUMN_SCORES[idx]
            print(f"✓ Column {idx+1} width correct (expected {EXPECTED_CM[idx]} cm, diff ±{round(diff/CM_TO_EMU,3)} cm) (+{COLUMN_SCORES[idx]})")
        else:
            print(f"✗ Column {idx+1} width incorrect – expected {EXPECTED_CM[idx]} cm, got {measured_cm[idx]} cm (diff {round(diff/CM_TO_EMU,3)} cm)")

    # 5) Final score (capped to 1.0)
    final_score = min(total_score, max_score)
    print(f"REWARD: {final_score}")
    return final_score

# -------------------------------------------------------------
# Execute verification when script is run
# -------------------------------------------------------------
if __name__ == "__main__":
    verify_table_column_widths(FILE_PATH)

