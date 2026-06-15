"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 48, I need the title to use the built-in “Outline” text effect, with the outline itself in solid black (#000000) and absolutely no fill (full transparency). How do I apply that in LibreOffice Impress?
Generated: 2025-09-11 00:24:31
Status: success
Model: azure-o3
Total Steps: 5
"""

import os
import zipfile
from pptx import Presentation
from lxml import etree


def verify_outline_title(file_path: str) -> float:
    """Verify that slide 48’s title uses the built-in *Outline* text effect
    with (1) *no fill* (fully transparent) and (2) a *solid black* (#000000)
    outline.  Progressive scoring awards points for each verified component.

    Returns a float between 0.0 and 1.0 and prints a detailed breakdown.
    """

    # XML namespace map for pptx drawing/presentation parts
    ns = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }

    max_score = 1.0      # maximum attainable score
    total_score = 0.0    # progressive score accumulator

    # ------------------------------------------------------------------
    # 0.  File existence & loadability  (prerequisite – NO POINTS!)
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ Presentation file not found.")
        print("REWARD: 0.0")
        return 0.0
    try:
        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        print(f"✓ Loaded presentation with {slide_count} slides")
    except Exception as e:
        print(f"✗ Unable to load presentation – {e}")
        print("REWARD: 0.0")
        return 0.0

    # ------------------------------------------------------------------
    # 1.  Verify slide 48 exists  (0.25 pts)
    # ------------------------------------------------------------------
    if slide_count >= 48:
        total_score += 0.25
        slide = prs.slides[47]  # zero-based index
        print("✓ Slide 48 exists (0.25 points)")
    else:
        print("✗ Presentation does not contain 48 slides")
        print(f"REWARD: {total_score}")
        return total_score

    # ------------------------------------------------------------------
    # 2.  Locate *title* placeholder on slide 48  (0.25 pts)
    # ------------------------------------------------------------------
    try:
        slide_partname = slide.part.partname  # e.g. /ppt/slides/slide48.xml
        with zipfile.ZipFile(file_path, "r") as z:
            slide_xml = z.read(slide_partname[1:])  # remove leading ‘/’
        root = etree.fromstring(slide_xml)

        # collect shapes whose placeholder type is title / ctrTitle
        title_shapes = [
            sp
            for sp in root.xpath(".//p:sp", namespaces=ns)
            if (sp.find(".//p:ph", namespaces=ns) is not None
                and sp.find(".//p:ph", namespaces=ns).get("type") in ("title", "ctrTitle"))
        ]

        if title_shapes:
            title_shape_xml = title_shapes[0]  # first title shape found
            total_score += 0.25
            print("✓ Found title placeholder on slide 48 (0.25 points)")
        else:
            print("✗ No title placeholder found on slide 48")
            print(f"REWARD: {total_score}")
            return total_score
    except Exception as e:
        print(f"✗ Error while locating title placeholder – {e}")
        print(f"REWARD: {total_score}")
        return total_score

    # ------------------------------------------------------------------
    # 3.  Check formatting of every run in the title
    #       a) *noFill*  →  full transparency               (0.25 pts)
    #       b) *ln/solidFill/srgbClr* = 000000 (black)     (0.25 pts)
    # ------------------------------------------------------------------
    runs = title_shape_xml.xpath(".//a:r", namespaces=ns)
    if not runs:
        print("✗ Title placeholder contains no text runs")
        print(f"REWARD: {total_score}")
        return total_score

    nofill_all = True
    outline_black_all = True

    for r in runs:
        rPr = r.find("./a:rPr", namespaces=ns)
        if rPr is None:
            nofill_all = False
            outline_black_all = False
            break

        # --- no fill? --------------------------------------------------
        if rPr.find("./a:noFill", namespaces=ns) is None:
            nofill_all = False

        # --- black outline? -------------------------------------------
        ln = rPr.find("./a:ln", namespaces=ns)
        ok_outline = False
        if ln is not None:
            srgb = ln.find("./a:solidFill/a:srgbClr", namespaces=ns)
            if srgb is not None and srgb.get("val", "").lower() == "000000":
                ok_outline = True
        if not ok_outline:
            outline_black_all = False

    # 3a. award points for *noFill*
    if nofill_all:
        total_score += 0.25
        print("✓ All title text runs have *no fill* (0.25 points)")
    else:
        print("✗ Not all title text runs are set to *no fill* (transparent)")

    # 3b. award points for solid black outline
    if outline_black_all:
        total_score += 0.25
        print("✓ All title text runs have a solid black outline (0.25 points)")
    else:
        print("✗ Not all title text runs have a solid black outline (#000000)")

    # ------------------------------------------------------------------
    final_score = min(total_score, max_score)
    print(f"\nTotal Score: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score


# ----------------------------------------------------------------------
# Execute verification when script is run directly
# ----------------------------------------------------------------------
if __name__ == "__main__":
    FILE = "/home/user/on_slide_48_i_need_the_title_to_use_the_built_in_outline_text_effect_with_the_outline_itself_in_soli_golden.pptx"
    verify_outline_title(FILE)
