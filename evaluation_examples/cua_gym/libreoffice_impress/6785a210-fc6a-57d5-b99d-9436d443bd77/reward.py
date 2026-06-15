"""
FINAL REWARD SCRIPT - SUCCESS
Task: Insert a footnote at the end of sentence 1 in paragraph 3 with text 'See Annex A.'
Generated: 2025-10-17 12:08:03
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation


def verify_footnote_task(file_path: str) -> float:
    """Verify that a footnote reference (superscript ¹) was inserted at the end
    of sentence 1 in paragraph 3 and that the corresponding footnote text
    "See Annex A." exists somewhere on the slide.

    Progressive scoring (max 1.0):
        0.4 – Superscript reference ¹ is present in any paragraph
        0.3 – The reference immediately follows a period ( … .¹ … )
        0.3 – A separate text element contains the footnote text
                "See Annex A." (case-insensitive, any leading ¹ allowed)
    """

    print(f"Verifying presentation: {file_path}\n")

    max_score = 1.0
    score = 0.0

    # --- Load presentation -------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File not found")
        return 0.0

    try:
        prs = Presentation(file_path)
        print(f"✓ File loaded – slide count: {len(prs.slides)}")
    except Exception as e:
        print(f"✗ Error loading PPTX: {e}")
        return 0.0

    # --- Requirement 1 : superscript reference exists ----------------------
    superscript_found = False
    positioned_correctly = False

    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or shape.text_frame is None:
                continue
            for paragraph in shape.text_frame.paragraphs:
                para_text = "".join(run.text for run in paragraph.runs)
                if "¹" in para_text:
                    superscript_found = True
                    print(f"✓ Found superscript reference in paragraph: '{para_text}'")
                    # Remove line-breaks to assess period immediately before superscript
                    if ".¹" in para_text.replace("\n", ""):
                        positioned_correctly = True
                    break
            if superscript_found:
                break
        if superscript_found:
            break

    if superscript_found:
        score += 0.4
        print("  → +0.4 points (reference exists)")
        if positioned_correctly:
            score += 0.3
            print("  → +0.3 points (reference placed immediately after period)")
        else:
            print("  → 0 points (reference misplaced – period not immediately before)")
    else:
        print("✗ No superscript reference found")

    # --- Requirement 2 : footnote text exists ------------------------------
    footnote_text_found = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            txt = shape.text.strip()
            if not txt:
                continue
            # Accept variants like '¹ See Annex A.' or 'See Annex A.'
            if "see annex a" in txt.lower():
                footnote_text_found = True
                print(f"✓ Found footnote text: '{txt}'")
                break
        if footnote_text_found:
            break

    if footnote_text_found:
        score += 0.3
        print("  → +0.3 points (footnote text correct)")
    else:
        print("✗ Footnote text 'See Annex A.' not found")

    # --- Final score -------------------------------------------------------
    final_score = min(score, max_score)
    print(f"\nREWARD: {final_score}")
    return final_score


# -------------------------------------------------------------------------
# Execute verification when script is run directly
if __name__ == "__main__":
    TEST_FILE = "/home/user/insert_a_footnote_at_the_end_of_sentence_1_in_paragraph_3_with_text_see_annex_a.pptx"
    verify_footnote_task(TEST_FILE)
