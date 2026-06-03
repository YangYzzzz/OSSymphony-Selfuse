"""
Reward Script: Extract presenter notes from two Impress files into a combined Word document.
Task ID: osworld_multi_apps_impress_notes_export_010
Domain: libreoffice_impress + libreoffice_writer (multi-app)

Scoring:
  Component 1: combined_notes.docx exists and is loadable                   (precondition gate)
  Component 2: 'Part 1' section header present with heading style            0.25
  Component 3: 'Part 2' section header present with heading style            0.25
  Component 4: All 5 Part1 notes appear in order before Part 2 section      0.30
  Component 5: All 7 Part2 notes appear in order after Part 2 header        0.20
  Total:                                                                      1.00
"""

import os
from docx import Document
from pptx import Presentation

WORKDIR = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_impress_notes_export_010'

PART1_PPTX = f'{WORKDIR}/Part1_Intro.pptx'
PART2_PPTX = f'{WORKDIR}/Part2_Advanced.pptx'
OUTPUT_DOCX = f'{WORKDIR}/combined_notes.docx'


def get_pptx_notes(pptx_path):
    """
    Extract per-slide presenter notes from a pptx file.
    Returns list of stripped note strings (one per slide), empty strings for blank notes.
    """
    prs = Presentation(pptx_path)
    notes_list = []
    for slide in prs.slides:
        try:
            text = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            text = ''
        notes_list.append(text)
    return notes_list


def texts_match_ordered(doc_paragraphs, expected_notes, start_idx):
    """
    Check that each note in expected_notes appears (in order) in doc_paragraphs
    starting from start_idx. Each expected note is verified by checking if the
    paragraph text starts with the same first 50 characters (allows for minor
    truncation differences). Returns (matched_count, next_idx).
    """
    idx = start_idx
    matched = 0
    for note in expected_notes:
        if not note:
            matched += 1
            continue
        prefix = note[:50]
        found = False
        search_end = min(idx + len(expected_notes) + 2, len(doc_paragraphs))
        for j in range(idx, search_end):
            if doc_paragraphs[j].text.strip().startswith(prefix):
                idx = j + 1
                found = True
                matched += 1
                break
        if not found:
            # Do not advance idx; stop searching for subsequent notes
            break
    return matched, idx


def verify_task(docx_path, part1_pptx, part2_pptx):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition gate: load combined_notes.docx
    try:
        doc = Document(docx_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load combined_notes.docx at {docx_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    paragraphs = doc.paragraphs

    # Load expected notes from the pptx files (on the VM)
    try:
        part1_notes = get_pptx_notes(part1_pptx)
        part2_notes = get_pptx_notes(part2_pptx)
    except Exception as e:
        print(f"CRITICAL: Cannot read source pptx files: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 2: 'Part 1' section header present with heading-like style (0.25 points)
    # This FAILS on initial_env (no docx exists) → PASSES on golden_env
    try:
        part1_header_idx = None
        for i, para in enumerate(paragraphs):
            text = para.text.strip()
            style_name = para.style.name.lower() if para.style else ''
            # Accept 'Part 1' or 'Part1' with a heading style
            if ('part 1' in text.lower() or 'part1' in text.lower()) and \
               ('heading' in style_name or 'title' in style_name):
                part1_header_idx = i
                break

        if part1_header_idx is not None:
            print(f"PASS: Component 2 — 'Part 1' heading found at paragraph {part1_header_idx} "
                  f"(style={paragraphs[part1_header_idx].style.name}) (0.25 pts)")
            total_score += 0.25
        else:
            print("FAIL: Component 2 — No 'Part 1' heading with heading style found in document")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: 'Part 2' section header present with heading-like style (0.25 points)
    # and appears AFTER 'Part 1' header (verifying order)
    try:
        part2_header_idx = None
        search_start = (part1_header_idx + 1) if part1_header_idx is not None else 0
        for i in range(search_start, len(paragraphs)):
            text = paragraphs[i].text.strip()
            style_name = paragraphs[i].style.name.lower() if paragraphs[i].style else ''
            if ('part 2' in text.lower() or 'part2' in text.lower()) and \
               ('heading' in style_name or 'title' in style_name):
                part2_header_idx = i
                break

        if part2_header_idx is not None:
            print(f"PASS: Component 3 — 'Part 2' heading found at paragraph {part2_header_idx} "
                  f"(style={paragraphs[part2_header_idx].style.name}), after 'Part 1' (0.25 pts)")
            total_score += 0.25
        else:
            print("FAIL: Component 3 — No 'Part 2' heading with heading style found after 'Part 1'")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: All 5 Part1 notes present in correct order before Part 2 header (0.30 points)
    # Progressive: award (matched/total) * 0.30
    try:
        if part1_header_idx is not None:
            start_idx = part1_header_idx + 1
            end_idx = part2_header_idx if part2_header_idx is not None else len(paragraphs)
            body_paras_part1 = paragraphs[start_idx:end_idx]

            matched = 0
            expected_count = len(part1_notes)
            note_idx_in_body = 0
            for note in part1_notes:
                if not note:
                    matched += 1
                    continue
                prefix = note[:50]
                found = False
                for j in range(note_idx_in_body, len(body_paras_part1)):
                    if body_paras_part1[j].text.strip().startswith(prefix):
                        note_idx_in_body = j + 1
                        found = True
                        matched += 1
                        break
                if not found:
                    print(f"  FAIL sub: Part1 slide note not found: {repr(prefix)}")

            score_comp4 = (matched / expected_count) * 0.30 if expected_count > 0 else 0.0
            if matched == expected_count:
                print(f"PASS: Component 4 — All {expected_count} Part1 notes present in order "
                      f"(0.30 pts)")
            else:
                print(f"PARTIAL: Component 4 — {matched}/{expected_count} Part1 notes found in order "
                      f"({score_comp4:.2f} pts)")
            total_score += score_comp4
        else:
            print("FAIL: Component 4 — Cannot evaluate Part1 notes (no 'Part 1' header found)")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: All 7 Part2 notes present in correct order after Part 2 header (0.20 points)
    # Progressive: award (matched/total) * 0.20
    try:
        if part2_header_idx is not None:
            start_idx = part2_header_idx + 1
            body_paras_part2 = paragraphs[start_idx:]

            matched = 0
            expected_count = len(part2_notes)
            note_idx_in_body = 0
            for note in part2_notes:
                if not note:
                    matched += 1
                    continue
                prefix = note[:50]
                found = False
                for j in range(note_idx_in_body, len(body_paras_part2)):
                    if body_paras_part2[j].text.strip().startswith(prefix):
                        note_idx_in_body = j + 1
                        found = True
                        matched += 1
                        break
                if not found:
                    print(f"  FAIL sub: Part2 slide note not found: {repr(prefix)}")

            score_comp5 = (matched / expected_count) * 0.20 if expected_count > 0 else 0.0
            if matched == expected_count:
                print(f"PASS: Component 5 — All {expected_count} Part2 notes present in order "
                      f"(0.20 pts)")
            else:
                print(f"PARTIAL: Component 5 — {matched}/{expected_count} Part2 notes found in order "
                      f"({score_comp5:.2f} pts)")
            total_score += score_comp5
        else:
            print("FAIL: Component 5 — Cannot evaluate Part2 notes (no 'Part 2' header found)")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = round(min(total_score, 1.0), 4)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Main entrypoint
if not os.path.exists(OUTPUT_DOCX):
    print(f"File not found: {OUTPUT_DOCX}")
    print("REWARD: 0.0")
else:
    verify_task(OUTPUT_DOCX, PART1_PPTX, PART2_PPTX)
