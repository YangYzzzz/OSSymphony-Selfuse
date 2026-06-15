"""
Initial Setup: Extract presenter notes from two presentations and merge into a Word doc
Task ID: osworld_multi_apps_impress_notes_export_010
Domain: libreoffice_impress (multi-app: Impress + Writer)

Creates:
  - /home/user/Desktop/Part1_Intro.pptx  (5 slides with speaker notes)
  - /home/user/Desktop/Part2_Advanced.pptx  (7 slides with speaker notes)

Opens Part1_Intro.pptx in LibreOffice Impress for the GUI agent.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

DESKTOP = '/home/user/Desktop'
PART1_PATH = f'{DESKTOP}/Part1_Intro.pptx'
PART2_PATH = f'{DESKTOP}/Part2_Advanced.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_part1():
    """Create Part1_Intro.pptx with 5 slides, each with speaker notes."""
    prs = Presentation()

    # Slide 1: Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Introduction to Machine Learning"
    slide1.placeholders[1].text = "A Comprehensive Overview"
    slide1.notes_slide.notes_text_frame.text = (
        "Welcome everyone to today's session on machine learning fundamentals. "
        "This introduction covers supervised, unsupervised, and reinforcement learning. "
        "Remind the audience to silence their phones before we begin."
    )

    # Slide 2: What is ML
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "What is Machine Learning?"
    slide2.placeholders[1].text = (
        "- Learning from data without explicit programming\n"
        "- Pattern recognition at scale\n"
        "- Iterative model improvement\n"
        "- Applications across every industry"
    )
    slide2.notes_slide.notes_text_frame.text = (
        "Emphasize the distinction between traditional programming and ML. "
        "In traditional programming, rules are explicitly coded. "
        "In ML, the model discovers patterns from training data. "
        "Give the Netflix recommendation system as a relatable example."
    )

    # Slide 3: Types of ML
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Types of Machine Learning"
    slide3.placeholders[1].text = (
        "Supervised Learning\n"
        "- Labeled training data\n"
        "- Classification and regression\n\n"
        "Unsupervised Learning\n"
        "- Unlabeled data, discover structure\n"
        "- Clustering and dimensionality reduction"
    )
    slide3.notes_slide.notes_text_frame.text = (
        "Walk through each category with concrete examples. "
        "Supervised: email spam detection uses labeled spam/not-spam emails. "
        "Unsupervised: customer segmentation groups buyers by behavior without predefined labels. "
        "Pause here for questions before moving to the next section."
    )

    # Slide 4: Key Algorithms
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Foundational Algorithms"
    slide4.placeholders[1].text = (
        "- Linear Regression: continuous value prediction\n"
        "- Decision Trees: rule-based classification\n"
        "- Random Forests: ensemble of decision trees\n"
        "- Neural Networks: deep learning foundation\n"
        "- Support Vector Machines: margin-based classification"
    )
    slide4.notes_slide.notes_text_frame.text = (
        "Do not go into mathematical depth here; keep it conceptual. "
        "Mention that decision trees are the easiest to explain to non-technical stakeholders. "
        "Neural networks will be covered in detail in Part 2. "
        "Reference the scikit-learn library for hands-on practice."
    )

    # Slide 5: Getting Started
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Getting Started with ML"
    slide5.placeholders[1].text = (
        "Step 1: Define the problem clearly\n"
        "Step 2: Collect and clean data\n"
        "Step 3: Choose the right algorithm\n"
        "Step 4: Train, validate, and iterate\n"
        "Step 5: Deploy and monitor"
    )
    slide5.notes_slide.notes_text_frame.text = (
        "Stress that data quality is more important than algorithm choice. "
        "A clean, well-labeled dataset with a simple algorithm often outperforms "
        "complex models on dirty data. "
        "Recommend Kaggle competitions for practical experience. "
        "This concludes Part 1; take a 10-minute break before Part 2 begins."
    )

    prs.save(PART1_PATH)
    print(f'Part1_Intro.pptx created: {PART1_PATH}')


def create_part2():
    """Create Part2_Advanced.pptx with 7 slides, each with speaker notes."""
    prs = Presentation()

    # Slide 1: Advanced Topics Overview
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Advanced Machine Learning Techniques"
    slide1.placeholders[1].text = "Deep Dives and Practical Applications"
    slide1.notes_slide.notes_text_frame.text = (
        "Welcome back from the break. "
        "Part 2 assumes familiarity with the fundamentals covered in Part 1. "
        "We will explore deep learning, transfer learning, and model optimization strategies."
    )

    # Slide 2: Deep Learning
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Deep Learning Architectures"
    slide2.placeholders[1].text = (
        "- Convolutional Neural Networks (CNN): image recognition\n"
        "- Recurrent Neural Networks (RNN): sequential data\n"
        "- Transformers: attention-based language models\n"
        "- Autoencoders: unsupervised representation learning"
    )
    slide2.notes_slide.notes_text_frame.text = (
        "Show the diagram of a CNN architecture if available. "
        "CNNs learn spatial hierarchies of features automatically. "
        "Transformers have largely replaced RNNs for NLP tasks since 2018. "
        "Mention BERT and GPT as landmark transformer models."
    )

    # Slide 3: Transfer Learning
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Transfer Learning"
    slide3.placeholders[1].text = (
        "Leverage pre-trained models for new tasks\n\n"
        "Benefits:\n"
        "- Reduced training time and compute\n"
        "- Better performance with limited data\n"
        "- Knowledge from large datasets transfers\n\n"
        "Popular base models: ResNet, VGG, BERT, GPT"
    )
    slide3.notes_slide.notes_text_frame.text = (
        "Transfer learning is arguably the most impactful practical technique today. "
        "A model trained on ImageNet (1.2M images) can be fine-tuned on just a few hundred images. "
        "Hugging Face is the go-to library for NLP transfer learning. "
        "Discuss domain adaptation when source and target domains differ significantly."
    )

    # Slide 4: Hyperparameter Tuning
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Hyperparameter Optimization"
    slide4.placeholders[1].text = (
        "Key hyperparameters to tune:\n"
        "- Learning rate (most critical)\n"
        "- Batch size\n"
        "- Number of layers and units\n"
        "- Regularization (dropout, L2)\n\n"
        "Methods: Grid Search, Random Search, Bayesian Optimization"
    )
    slide4.notes_slide.notes_text_frame.text = (
        "The learning rate is the single most important hyperparameter. "
        "Start with learning rate range test (LR finder) before full training. "
        "Bayesian optimization is more sample-efficient than grid or random search. "
        "Tools: Optuna, Ray Tune, Weights and Biases sweeps."
    )

    # Slide 5: Model Evaluation
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Rigorous Model Evaluation"
    slide5.placeholders[1].text = (
        "Beyond accuracy:\n"
        "- Precision, Recall, F1 for imbalanced classes\n"
        "- ROC-AUC for ranking ability\n"
        "- Confusion matrix for error analysis\n"
        "- Cross-validation for robust estimates\n\n"
        "Avoid data leakage at all costs"
    )
    slide5.notes_slide.notes_text_frame.text = (
        "Accuracy is misleading on imbalanced datasets. "
        "Example: 99% accuracy on fraud detection could mean simply predicting non-fraud every time. "
        "Always stratify your train/test split for classification problems. "
        "Data leakage is the most common mistake in competitions and real projects alike."
    )

    # Slide 6: MLOps and Deployment
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "MLOps: From Research to Production"
    slide6.placeholders[1].text = (
        "Production ML challenges:\n"
        "- Model versioning and reproducibility\n"
        "- Continuous training pipelines\n"
        "- Data and concept drift monitoring\n"
        "- A/B testing and canary releases\n\n"
        "Tools: MLflow, Kubeflow, SageMaker"
    )
    slide6.notes_slide.notes_text_frame.text = (
        "The gap between research prototype and production system is often underestimated. "
        "Concept drift occurs when real-world data distribution shifts after deployment. "
        "Mention the 2020 model degradation during COVID as a real-world drift example. "
        "Shadow deployment is a low-risk strategy: run new model alongside old model before switching."
    )

    # Slide 7: Future Directions
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Future Directions in ML"
    slide7.placeholders[1].text = (
        "Emerging research areas:\n"
        "- Foundation models and scaling laws\n"
        "- Federated learning for privacy\n"
        "- Neurosymbolic AI\n"
        "- Efficient inference and edge deployment\n"
        "- Multimodal learning"
    )
    slide7.notes_slide.notes_text_frame.text = (
        "The field is moving extraordinarily fast; encourage continuous learning. "
        "Foundation models like GPT-4 demonstrate emergent capabilities at scale. "
        "Federated learning is critical for healthcare where data cannot leave hospitals. "
        "Close with Q&A session. Thank the audience for their attention and participation."
    )

    prs.save(PART2_PATH)
    print(f'Part2_Advanced.pptx created: {PART2_PATH}')


def main():
    os.makedirs(DESKTOP, exist_ok=True)

    create_part1()
    create_part2()

    # Ensure combined_notes.docx does NOT exist in initial state
    combined_path = f'{DESKTOP}/combined_notes.docx'
    if os.path.exists(combined_path):
        os.remove(combined_path)
        print(f'Removed pre-existing combined_notes.docx')

    # GUI-ready startup: open Part1_Intro.pptx in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{PART1_PATH}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with Part1_Intro.pptx (DISPLAY=:0)')


main()
