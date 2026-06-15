"""
Initial Setup: Eco Report presentation with a table on slide 1
Task ID: impress_tct_008
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_008'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title + Table ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Title text box
    txBox = slide1.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Annual Sustainability Performance Report 2025"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x7D, 0x32)

    # 5-col x 4-row table
    headers = ["Category", "Q1 Target", "Q1 Actual", "Q2 Target", "Q2 Actual"]
    data = [
        ["Carbon Emissions (tons CO₂)", "1,200", "1,145", "1,150", "1,098"],
        ["Water Usage (m³)", "8,500", "8,320", "8,200", "7,950"],
        ["Renewable Energy (%)", "45%", "47.2%", "50%", "52.1%"],
    ]

    tbl_shape = slide1.shapes.add_table(
        rows=4, cols=5,
        left=Inches(0.8), top=Inches(2.0),
        width=Inches(11.5), height=Inches(3.0),
    )
    table = tbl_shape.table

    # Set column widths
    col_widths = [Inches(3.2), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header row — 12pt, black, NOT bold, no background
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = hdr
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.name = "Calibri"
            r.font.size = Pt(12)
            r.font.bold = False
            r.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        # Explicitly no fill on header cells
        cell.fill.background()

    # Data rows — 11pt, dark gray
    for ri, row_data in enumerate(data, start=1):
        for ci, val in enumerate(row_data):
            cell = table.cell(ri, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            for r in p.runs:
                r.font.name = "Calibri"
                r.font.size = Pt(11)
                r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            # Light gray alternating rows
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
            else:
                cell.fill.background()

    # ── Slide 2: Key Achievements ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.9))
    p2 = txBox2.text_frame.paragraphs[0]
    p2.text = "Key Achievements — First Half 2025"
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.runs[0]
    r2.font.name = "Calibri"
    r2.font.size = Pt(24)
    r2.font.bold = True
    r2.font.color.rgb = RGBColor(0x2E, 0x7D, 0x32)

    body2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(10.5), Inches(5.0))
    tf2 = body2.text_frame
    tf2.word_wrap = True
    bullets = [
        "Reduced carbon footprint by 8.4% against the annual baseline",
        "Achieved 52.1% renewable energy adoption ahead of schedule",
        "Installed 340 kW rooftop solar array at the Phoenix facility",
        "Eliminated single-use plastics from all cafeteria operations",
        "Partnered with local watershed council for water reclamation",
    ]
    for i, txt in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = txt
        p.space_after = Pt(10)
        for r in p.runs:
            r.font.name = "Calibri"
            r.font.size = Pt(16)
            r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ── Slide 3: Strategic Outlook ──
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.9))
    p3 = txBox3.text_frame.paragraphs[0]
    p3.text = "Strategic Outlook — 2025-2027"
    p3.alignment = PP_ALIGN.LEFT
    r3 = p3.runs[0]
    r3.font.name = "Calibri"
    r3.font.size = Pt(24)
    r3.font.bold = True
    r3.font.color.rgb = RGBColor(0x2E, 0x7D, 0x32)

    body3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(10.5), Inches(5.0))
    tf3 = body3.text_frame
    tf3.word_wrap = True
    outlook_items = [
        "Achieve net-zero Scope 1 & 2 emissions by Q4 2026",
        "Expand EV fleet to cover 75% of company logistics vehicles",
        "Deploy AI-driven energy management across all 12 facilities",
        "Establish biodiversity offset program with 500-acre reforestation pledge",
    ]
    for i, txt in enumerate(outlook_items):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = txt
        p.space_after = Pt(10)
        for r in p.runs:
            r.font.name = "Calibri"
            r.font.size = Pt(16)
            r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ── Slide 4: Thank You / Contact ──
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    fill4 = slide4.background.fill
    fill4.solid()
    fill4.fore_color.rgb = RGBColor(0x2E, 0x7D, 0x32)

    txBox4 = slide4.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2.5))
    tf4 = txBox4.text_frame
    tf4.word_wrap = True
    p4 = tf4.paragraphs[0]
    p4.text = "Thank You"
    p4.alignment = PP_ALIGN.CENTER
    r4 = p4.runs[0]
    r4.font.name = "Calibri"
    r4.font.size = Pt(40)
    r4.font.bold = True
    r4.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    p4b = tf4.add_paragraph()
    p4b.text = "sustainability@greenearth-corp.com  |  www.greenearth-corp.com"
    p4b.alignment = PP_ALIGN.CENTER
    r4b = p4b.runs[0]
    r4b.font.name = "Calibri"
    r4b.font.size = Pt(18)
    r4b.font.color.rgb = RGBColor(0xE8, 0xF5, 0xE9)

    # Save
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
