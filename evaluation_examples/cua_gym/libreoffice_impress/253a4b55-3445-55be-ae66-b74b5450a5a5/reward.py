"""
Reward Script: Make header row text bold, 16pt, white (#FFFFFF) with dark green (#1B5E20) background
Task ID: impress_tct_008
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): All header cells have bold text
  Component 2 (0.25): All header cells have 16pt (203200 EMU) font size
  Component 3 (0.25): All header cells have white (#FFFFFF) text color
  Component 4 (0.25): All header cells have dark green (#1B5E20) background fill
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_008'


def persist_app_state(domain):
    """Save any unsaved GUI edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print(f"PERSIST: ctrl+s sent for {domain}")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def find_table(slide):
    """Find the first table shape on the slide."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return shape.table
    return None


def get_cell_fill_color(cell):
    """Extract solid fill color from a table cell via XML."""
    tcPr = cell._tc.find(qn('a:tcPr'))
    if tcPr is None:
        return None
    solidFill = tcPr.find(qn('a:solidFill'))
    if solidFill is None:
        return None
    srgb = solidFill.find(qn('a:srgbClr'))
    if srgb is not None:
        return srgb.get('val')
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: slide 1 exists and has a table
    if len(prs.slides) < 1:
        print("FAIL: No slides in presentation")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]
    table = find_table(slide)
    if table is None:
        print("FAIL: No table found on slide 1")
        print("REWARD: 0.0")
        return 0.0

    num_cols = len(table.columns)
    if num_cols == 0:
        print("FAIL: Table has 0 columns")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Found table with {len(table.rows)} rows x {num_cols} cols")

    # Component 1: All header cells have bold text (0.25 points)
    try:
        bold_count = 0
        for ci in range(num_cols):
            cell = table.cell(0, ci)
            all_bold = True
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        if run.font.bold is not True:
                            all_bold = False
            if all_bold:
                bold_count += 1
        if bold_count == num_cols:
            print(f"PASS: Component 1 -- All {num_cols} header cells have bold text (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 -- Only {bold_count}/{num_cols} header cells have bold text")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: All header cells have 16pt font size (203200 EMU) (0.25 points)
    try:
        size_count = 0
        expected_size = 203200  # 16pt in EMU
        for ci in range(num_cols):
            cell = table.cell(0, ci)
            correct_size = True
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        if run.font.size != expected_size:
                            correct_size = False
                            print(f"  DEBUG: Cell[0,{ci}] run size={run.font.size}, expected={expected_size}")
            if correct_size:
                size_count += 1
        if size_count == num_cols:
            print(f"PASS: Component 2 -- All {num_cols} header cells have 16pt font (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 -- Only {size_count}/{num_cols} header cells have 16pt font")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: All header cells have white (#FFFFFF) text color (0.25 points)
    try:
        color_count = 0
        for ci in range(num_cols):
            cell = table.cell(0, ci)
            correct_color = True
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        try:
                            if run.font.color.type is not None:
                                rgb_str = str(run.font.color.rgb).upper()
                                if rgb_str != 'FFFFFF':
                                    correct_color = False
                                    print(f"  DEBUG: Cell[0,{ci}] run color={rgb_str}, expected=FFFFFF")
                            else:
                                correct_color = False
                                print(f"  DEBUG: Cell[0,{ci}] run color type is None (inherited)")
                        except Exception:
                            correct_color = False
            if correct_color:
                color_count += 1
        if color_count == num_cols:
            print(f"PASS: Component 3 -- All {num_cols} header cells have white text (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 -- Only {color_count}/{num_cols} header cells have white text")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: All header cells have dark green (#1B5E20) background fill (0.25 points)
    try:
        fill_count = 0
        for ci in range(num_cols):
            cell = table.cell(0, ci)
            fill_color = get_cell_fill_color(cell)
            if fill_color is not None and fill_color.upper() == '1B5E20':
                fill_count += 1
            else:
                print(f"  DEBUG: Cell[0,{ci}] fill={fill_color}, expected=1B5E20")
        if fill_count == num_cols:
            print(f"PASS: Component 4 -- All {num_cols} header cells have #1B5E20 background (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 4 -- Only {fill_count}/{num_cols} header cells have #1B5E20 background")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist any unsaved GUI state before verification
persist_app_state("libreoffice_impress")

# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
