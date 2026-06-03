"""
Initial Setup: Create a 15-slide presentation with no custom slideshows
Task ID: impress_ndo_088
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_088'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, point in enumerate(bullet_points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
    return slide


def add_blank_with_text(prs, title_text, body_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    # Title textbox
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
    p.alignment = PP_ALIGN.LEFT
    # Body textbox
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = body_text
    p2.font.size = Pt(16)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide (shared)
    add_title_slide(prs, "Versatile Platform Architecture",
                    "Cross-Functional Strategy & Technical Deep Dive\nQ2 2025 Review")

    # Slide 2: Business - Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "Revenue grew 23% year-over-year to $47.8M",
        "Customer acquisition cost decreased by 15%",
        "Net promoter score improved from 62 to 74",
        "Three new enterprise clients onboarded in Q1",
        "Operating margin expanded to 18.3%"
    ])

    # Slide 3: Business - Market Overview
    add_content_slide(prs, "Market Landscape & Positioning", [
        "Total addressable market expanded to $12.4B",
        "Competitive differentiation through AI-driven analytics",
        "Strategic partnerships with Salesforce and AWS",
        "APAC region showing 40% growth trajectory",
        "Enterprise segment represents 65% of new bookings"
    ])

    # Slide 4: Technical - System Architecture
    add_content_slide(prs, "System Architecture Overview", [
        "Microservices deployed across 3 AWS regions",
        "Event-driven architecture using Apache Kafka",
        "GraphQL API gateway with rate limiting",
        "Redis caching layer with 99.7% hit rate",
        "PostgreSQL with read replicas for analytics workloads"
    ])

    # Slide 5: Technical - Infrastructure
    add_content_slide(prs, "Infrastructure & DevOps Pipeline", [
        "Kubernetes clusters managing 240+ pods",
        "CI/CD pipeline averaging 12-minute deployments",
        "Infrastructure as Code via Terraform modules",
        "Prometheus + Grafana monitoring stack",
        "Automated canary deployments with 1% traffic rollout"
    ])

    # Slide 6: Technical - Security
    add_content_slide(prs, "Security Architecture & Compliance", [
        "Zero-trust network architecture implemented",
        "SOC 2 Type II certification renewed March 2025",
        "End-to-end encryption with AES-256-GCM",
        "Automated vulnerability scanning in CI pipeline",
        "Mean time to remediate critical CVEs: 4.2 hours"
    ])

    # Slide 7: Technical - Performance
    add_content_slide(prs, "Performance Engineering Results", [
        "P99 latency reduced from 340ms to 89ms",
        "Database query optimization saved 60% compute costs",
        "CDN cache hit ratio improved to 94.2%",
        "WebSocket connections handling 50K concurrent users",
        "Load testing validated 3x current peak capacity"
    ])

    # Slide 8: Business - Financial Projections
    add_content_slide(prs, "Financial Projections FY2025-2026", [
        "Projected ARR of $62M by end of FY2025",
        "Gross margin target: 78% (current: 74%)",
        "R&D investment increasing to 28% of revenue",
        "Customer lifetime value: $185K (up from $142K)",
        "Path to profitability projected for Q3 2026"
    ])

    # Slide 9: Business - Go-to-Market Strategy
    add_content_slide(prs, "Go-to-Market Strategy Update", [
        "Partner channel contributing 30% of pipeline",
        "Product-led growth motion launched in January",
        "Self-serve tier conversion rate at 8.4%",
        "Enterprise sales cycle reduced to 67 days average",
        "Marketing qualified leads up 45% quarter-over-quarter"
    ])

    # Slide 10: Business - Customer Success
    add_content_slide(prs, "Customer Success & Retention Metrics", [
        "Net revenue retention rate: 118%",
        "Churn reduced to 3.2% annually (from 5.1%)",
        "Customer health score averaging 82/100",
        "Onboarding time decreased from 21 to 9 days",
        "Support ticket resolution time: 2.4 hours average"
    ])

    # Slide 11: Mixed - Product Roadmap
    add_content_slide(prs, "Product Roadmap Highlights", [
        "AI copilot feature entering beta in June 2025",
        "Mobile app redesign shipping Q3 2025",
        "Advanced analytics dashboard for enterprise tier",
        "API v3 with improved developer experience",
        "Internationalization: 8 new languages by year-end"
    ])

    # Slide 12: Technical - Data Platform
    add_content_slide(prs, "Data Platform & ML Pipeline", [
        "Data lake processing 2.3TB daily ingestion",
        "Feature store serving 150+ ML models in production",
        "Real-time inference pipeline with sub-10ms latency",
        "A/B testing framework running 40+ experiments monthly",
        "Data quality score improved to 97.8% accuracy"
    ])

    # Slide 13: Mixed - Team & Culture
    add_content_slide(prs, "Team Growth & Engineering Culture", [
        "Engineering team grew from 45 to 72 members",
        "4 new senior staff engineers hired from FAANG",
        "Internal tech talks program with 120+ attendees",
        "Hackathon yielded 3 features now in production",
        "Employee satisfaction score: 4.6/5.0"
    ])

    # Slide 14: Mixed - Risk Assessment
    add_content_slide(prs, "Risk Assessment & Mitigation", [
        "Supply chain dependencies mapped and monitored",
        "Disaster recovery RTO reduced to 15 minutes",
        "Regulatory compliance across 12 jurisdictions",
        "Vendor concentration risk addressed with multi-cloud",
        "Talent retention program showing 92% effectiveness"
    ])

    # Slide 15: Closing Slide (shared)
    add_title_slide(prs, "Next Steps & Action Items",
                    "Technical Review: Architecture Committee Meeting — May 12\n"
                    "Business Review: Board Presentation — May 20\n"
                    "Contact: strategy@versatileplatform.com")

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
