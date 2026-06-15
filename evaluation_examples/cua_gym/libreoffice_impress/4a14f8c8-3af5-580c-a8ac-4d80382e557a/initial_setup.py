"""
Initial Setup: Create Flow_Diagram.pptx with 4 arrow shapes on slide 8, no animations
Task ID: impress_fix_085
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_085'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Process Flow Diagram"
    slide1.placeholders[1].text = "Q2 2025 Operations Review\nPrepared by: Elena Vasquez, Operations Director"

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Project Overview"
    tf = slide2.placeholders[1].text_frame
    tf.text = "This presentation outlines the end-to-end process flow for our manufacturing pipeline."
    p = tf.add_paragraph()
    p.text = "Key areas covered: procurement, assembly, quality control, and distribution."
    p = tf.add_paragraph()
    p.text = "Timeline: April - September 2025"

    # --- Slide 3: Procurement Phase ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Phase 1: Procurement"
    tf = slide3.placeholders[1].text_frame
    tf.text = "Vendor selection and raw material sourcing"
    for item in ["RFQ distribution to 12 approved vendors",
                 "Material quality certification review",
                 "Purchase order generation and approval workflow",
                 "Delivery scheduling with 3-week lead time"]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 4: Assembly Phase ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Phase 2: Assembly Line Setup"
    tf = slide4.placeholders[1].text_frame
    tf.text = "Three parallel assembly lines operating 16 hours/day"
    for item in ["Line A: Electronic component integration",
                 "Line B: Mechanical housing assembly",
                 "Line C: Final product integration and sealing",
                 "Estimated throughput: 2,400 units per shift"]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 5: Quality Control ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Phase 3: Quality Control"
    tf = slide5.placeholders[1].text_frame
    tf.text = "Multi-stage inspection protocol"
    for item in ["Visual inspection station (automated camera system)",
                 "Electrical continuity and performance testing",
                 "Environmental stress screening (thermal cycling)",
                 "Final audit sampling at 5% rate"]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 6: Distribution ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Phase 4: Distribution Network"
    tf = slide6.placeholders[1].text_frame
    tf.text = "Regional warehouse allocation strategy"
    for item in ["Northeast hub: Edison, NJ (35% of volume)",
                 "Southeast hub: Atlanta, GA (25% of volume)",
                 "Central hub: Dallas, TX (20% of volume)",
                 "West Coast hub: Ontario, CA (20% of volume)"]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 7: Timeline ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Implementation Timeline"
    tf = slide7.placeholders[1].text_frame
    tf.text = "Key milestones and deliverables"
    for item in ["April 2025: Vendor contracts finalized",
                 "May 2025: Assembly line commissioning",
                 "June 2025: First production run",
                 "July 2025: Full-scale production begins",
                 "September 2025: Distribution network fully operational"]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 8: Process Flow Arrows (THE KEY SLIDE) ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add a title text box
    title_box = slide8.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Process Flow Animation"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # Labels for the 4 arrows
    labels = ["Procurement", "Assembly", "Quality Control", "Distribution"]
    colors = [
        RGBColor(0x2E, 0x86, 0xC1),  # Blue
        RGBColor(0x28, 0xB4, 0x63),  # Green
        RGBColor(0xE6, 0x7E, 0x22),  # Orange
        RGBColor(0x8E, 0x44, 0xAD),  # Purple
    ]
    y_positions = [Inches(1.5), Inches(3.0), Inches(4.5), Inches(6.0)]

    for i in range(4):
        # Add right arrow shape at left side (x=0.5in)
        arrow = slide8.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(0.5),        # left position
            y_positions[i],      # vertical position
            Inches(2.5),         # width
            Inches(0.8),         # height
        )
        # Fill color
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = colors[i]
        # Border
        arrow.line.color.rgb = colors[i]
        arrow.line.width = Pt(1)

        # Add label text inside arrow
        tf = arrow.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = labels[i]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Add a descriptive label to the right of initial position
        label_box = slide8.shapes.add_textbox(
            Inches(3.2), y_positions[i] + Inches(0.15),
            Inches(3), Inches(0.5)
        )
        ltf = label_box.text_frame
        lp = ltf.paragraphs[0]
        lp.text = f"Phase {i+1}: {labels[i]}"
        lrun = lp.runs[0]
        lrun.font.size = Pt(11)
        lrun.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        lrun.font.italic = True

    # --- Slide 9: Summary ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Summary & Next Steps"
    tf = slide9.placeholders[1].text_frame
    tf.text = "The process flow has been optimized for maximum efficiency."
    p = tf.add_paragraph()
    p.text = "Next review meeting: May 15, 2025"
    p = tf.add_paragraph()
    p.text = "Contact: elena.vasquez@company.com"

    # --- Slide 10: Thank You ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[0])
    slide10.shapes.title.text = "Thank You"
    slide10.placeholders[1].text = "Questions & Discussion"

    # NO animations applied anywhere - this is the initial state
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
