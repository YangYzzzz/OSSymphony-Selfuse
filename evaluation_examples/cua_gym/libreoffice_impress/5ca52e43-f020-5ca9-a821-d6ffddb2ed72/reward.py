"""
FINAL REWARD SCRIPT - SUCCESS
Task: Center-align the document title (first paragraph) and set spacing after 12 pt.
Generated: 2025-10-17 08:02:21
Status: success
Model: azure-o3
Total Steps: 2
"""

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
import os

def verify_center_align_and_spacing(file_path: str) -> float:
    """Verify that the first paragraph of the first slide is
    1) centre-aligned and
    2) has 12 pt space-after.

    Returns a progressive score between 0.0 and 1.0.
    """
    print(f"Verifying presentation file: {file_path}")
    score = 0.0
    max_score = 1.0

    # ---------- Load the file ----------
    if not os.path.exists(file_path):
        print("✗ File does not exist")
        return 0.0

    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded PPTX with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Failed to load PPTX: {e}")
        return 0.0

    if not prs.slides:
        print("✗ Presentation contains no slides")
        return 0.0

    # ---------- Locate the first non-empty paragraph on slide 1 ----------
    first_slide = prs.slides[0]
    first_paragraph = None

    for shape in first_slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        tf = shape.text_frame
        for para in tf.paragraphs:
            text = "".join(run.text for run in para.runs).strip()
            if text:  # first non-empty paragraph
                first_paragraph = para
                print(f"✓ Found first non-empty paragraph: '{text[:50]}'")
                break
        if first_paragraph is not None:
            break

    if first_paragraph is None:
        print("✗ Could not find any paragraph with text on the first slide")
        return 0.0

    # ---------- Check centre alignment ----------
    alignment = first_paragraph.alignment
    print("Paragraph alignment value:", alignment)
    if alignment == PP_ALIGN.CENTER:
        print("✓ Paragraph is centre-aligned (0.5 points)")
        score += 0.5
    else:
        print("✗ Paragraph is not centre-aligned")

    # ---------- Check spacing-after (12 pt) ----------
    space_after = first_paragraph.space_after  # Length object or None
    if space_after is not None:
        space_after_pt = space_after.pt  # convert to points
        print(f"Spacing after: {space_after_pt:.2f} pt")
        if abs(space_after_pt - 12.0) <= 0.2:  # allow tiny tolerance
            print("✓ Spacing after is 12 pt (0.5 points)")
            score += 0.5
        else:
            print("✗ Spacing after is not 12 pt")
    else:
        print("✗ Spacing after is not set")

    # ---------- Final score ----------
    final_score = min(score, max_score)
    print(f"Total score: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score

# --------------- Script entry point ---------------
if __name__ == "__main__":
    verify_center_align_and_spacing(
        "/home/user/center_align_the_document_title_first_paragraph_and_set_spacing_after_12_pt.pptx"
    )
