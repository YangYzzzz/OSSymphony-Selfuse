"""
FINAL REWARD SCRIPT - SUCCESS
Task: Insert header text 'Confidential' left-aligned on all pages.
Generated: 2025-10-17 12:42:35
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

"""
Reward Script
Task:  Insert header text 'Confidential' left-aligned on ALL slides of the presentation.
This script verifies two independent aspects on every slide:
1.  Coverage  – the text 'Confidential' is present on the slide (case-insensitive)
2.  Alignment – that header is LEFT-ALIGNED both in paragraph formatting AND positioned
                 near the left edge of the slide (shape.left below a threshold)

Scoring (progressive – 0.0‥1.0):
•  0.5 points are awarded proportionally to how many slides contain the header text
•  0.5 points are awarded proportionally to how many slides contain the header text AND
   meet the left-alignment criteria
The final score is the sum of the two parts, capped at 1.0.

Notes / Anti-hacking compliance:
•  No points for simply loading a file – prerequisite only
•  Verification uses real pptx parsing (python-pptx)
•  No subprocess usage, no hard-coded truth values
"""

def verify_confidential_header(file_path: str) -> float:
    max_score = 1.0
    print(f"Verifying file: {file_path}")

    # Basic existence check (no points for this):
    if not os.path.exists(file_path):
        print("✗ File does not exist")
        print("REWARD: 0.0")
        return 0.0

    # Attempt to load presentation (no points for this alone):
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Could not open presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    total_slides = len(prs.slides)
    if total_slides == 0:
        print("✗ Presentation has no slides")
        print("REWARD: 0.0")
        return 0.0

    print(f"✓ Loaded presentation with {total_slides} slides")

    # Counters for verification results
    slides_with_confidential = 0  # header text present
    slides_with_left_aligned_confidential = 0  # header text present AND verified left-aligned

    # Threshold for considering a shape to be near the left edge (EMU units)
    # 914400 EMU = 1 inch  → 400000 ≈ 0.44 inch (about 1.1 cm)
    LEFT_POS_THRESHOLD = 400_000

    # Iterate through slides and inspect shapes
    for idx, slide in enumerate(prs.slides, start=1):
        header_found = False
        header_left_aligned = False

        for shape in slide.shapes:
            # Need a text frame to contain paragraphs
            if not getattr(shape, 'has_text_frame', False):
                continue

            text = (shape.text or '').strip()
            if not text:
                continue

            # Check for the keyword 'Confidential' (case-insensitive)
            if 'confidential' in text.lower():
                header_found = True

                # 1) Paragraph alignment check – all paragraphs must be left or None (defaults to left)
                if all((p.alignment in (None, PP_PARAGRAPH_ALIGNMENT.LEFT)) for p in shape.text_frame.paragraphs):
                    # 2) Geometric position check – shape must be very close to the left edge
                    if shape.left <= LEFT_POS_THRESHOLD:
                        header_left_aligned = True
                break  # We only need to validate the first header found

        # Update counters after examining the slide
        if header_found:
            slides_with_confidential += 1
        if header_left_aligned:
            slides_with_left_aligned_confidential += 1

        print(f"Slide {idx}: header_found={header_found}, left_aligned={header_left_aligned}")

    # --- Progressive Scoring ---
    coverage_fraction = slides_with_confidential / total_slides
    alignment_fraction = slides_with_left_aligned_confidential / total_slides

    # 50% weight for each aspect
    coverage_score = 0.5 * coverage_fraction
    alignment_score = 0.5 * alignment_fraction

    total_score = round(min(coverage_score + alignment_score, max_score), 2)

    print(f"Coverage:  {slides_with_confidential}/{total_slides} → {coverage_score:.2f}")
    print(f"Alignment: {slides_with_left_aligned_confidential}/{total_slides} → {alignment_score:.2f}")
    print(f"Total score: {total_score}")
    print(f"REWARD: {total_score}")

    return total_score


if __name__ == "__main__":
    # Path provided by the task context – adjust if task specifies a different location
    FILE_PATH = "/home/user/insert_header_text_confidential_left_aligned_on_all_pages.pptx"
    verify_confidential_header(FILE_PATH)
