"""
FINAL REWARD SCRIPT - SUCCESS
Task: My deck has grown to 300+ slides, and I just noticed slide 244 is still on the default title layout. I’d like to swap only that slide over to the “Title and Two Content” layout in LibreOffice Impress so I can show a chart beside the bullet list. What’s the quickest way to make that change?
Generated: 2025-09-10 20:06:56
Status: success
Model: azure-o3
Total Steps: 6
"""

import os
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

def verify_slide_244_layout(presentation):
    """Verify that slide 244 (index 243) uses a Title and Two Content layout.
    Returns True when the slide layout name matches the expected layout *and*
    the placeholder structure contains at least 1 title and 2 body/object placeholders.
    """
    target_index = 243  # zero-based index for slide 244
    if target_index >= len(presentation.slides):
        print("✗ Presentation does not contain slide 244")
        return False

    slide = presentation.slides[target_index]

    # 1. Check layout name contains the key phrase (LibreOffice exports often call it "Two Content")
    layout_name = (slide.slide_layout.name or "").lower()
    name_ok = "two content" in layout_name  # covers "Title and Two Content" & "Two Content"
    print(f"Slide 244 layout name: '{slide.slide_layout.name}' – name check: {name_ok}")

    # 2. Check placeholder structure: ≥1 title + ≥2 body/object placeholders
    placeholder_types = [sh.placeholder_format.type for sh in slide.shapes if getattr(sh, 'is_placeholder', False)]
    title_count = sum(1 for p in placeholder_types if p == PP_PLACEHOLDER.TITLE)
    body_like_count = sum(1 for p in placeholder_types if p in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT))
    structure_ok = title_count >= 1 and body_like_count >= 2
    print(f"  Placeholder counts – title: {title_count}, body/object: {body_like_count} – structure check: {structure_ok}")

    return name_ok and structure_ok


def verify_task(file_path):
    """Reward-scoring verifier for the LibreOffice Impress task.

    Scoring rubric:
      • 0.85 points – Slide 244 correctly converted to "Title and Two Content" layout
      • 0.15 points – Deck still contains at least 300 slides (integrity check)
    Returns a float between 0.0 and 1.0.
    """
    print(f"Starting verification for: {file_path}\n")
    MAX_SCORE = 1.0
    score = 0.0

    # Prerequisite: file must exist and be loadable (no score awarded for simply loading).
    if not os.path.exists(file_path):
        print("✗ File not found – cannot verify task")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        print(f"✓ Presentation loaded – slide count: {slide_count}")
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # 1) Integrity check – deck still ≥ 300 slides (0.15 pts)
    if slide_count >= 300:
        print("✓ Slide count integrity check passed (≥ 300 slides)")
        score += 0.15
    else:
        print(f"✗ Slide count integrity check failed (found {slide_count}, need ≥ 300)")

    # 2) Core requirement – slide 244 layout corrected (0.85 pts)
    if verify_slide_244_layout(prs):
        print("✓ Slide 244 uses the correct 'Title and Two Content' layout")
        score += 0.85
    else:
        print("✗ Slide 244 layout verification failed")

    # Final score (cap at 1.0)
    final_score = min(score, MAX_SCORE)
    print(f"\nTotal score: {final_score} / {MAX_SCORE}")
    print(f"REWARD: {final_score}")
    return final_score

# -----------------------
# Execute when run as script
# -----------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/my_deck_has_grown_to_300_slides_and_i_just_noticed_slide_244_is_still_on_the_default_title_layout_id_golden.pptx"
    verify_task(FILE_PATH)

