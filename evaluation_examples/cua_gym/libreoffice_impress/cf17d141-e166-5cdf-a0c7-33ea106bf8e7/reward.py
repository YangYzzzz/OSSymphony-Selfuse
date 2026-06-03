"""
Reward Script: Add 100% stacked bar chart on slide 5 with revenue mix data
Task ID: impress_exec_089
Domain: libreoffice_impress
Scoring:
  - Component 1: Chart exists on slide 5 (0.20)
  - Component 2: Chart type is BAR_STACKED_100 (0.15)
  - Component 3: Chart title matches expected (0.15)
  - Component 4: Chart data (3 series, 4 categories, correct values) (0.25)
  - Component 5: Chart series colors match specification (0.25)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_089'


def persist_app_state(domain):
    """Save any unsaved LibreOffice edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for", domain)
    except Exception as e:
        print("PERSIST_WARN: save hook failed:", e)


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print("CRITICAL: Cannot load file {}: {}".format(file_path, e))
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 5 slides
    if len(prs.slides) < 5:
        print("FAIL: Presentation has only {} slides, need at least 5".format(len(prs.slides)))
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # slide 5 (0-indexed)

    # Find chart shape on slide 5
    chart_shape = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            chart_shape = shape
            break

    # Component 1: Chart exists on slide 5 (0.20 points)
    try:
        if chart_shape is not None:
            print("PASS: Component 1 -- Chart exists on slide 5 (0.20 pts)")
            total_score += 0.20
        else:
            print("FAIL: Component 1 -- No chart found on slide 5")
            # Without a chart, no further checks can pass
            final_score = min(total_score, 1.0)
            print("\nScore: {}/1.0".format(total_score))
            print("REWARD: {}".format(final_score))
            return final_score
    except Exception as e:
        print("ERROR: Component 1 -- {}".format(e))

    chart = chart_shape.chart

    # Component 2: Chart type is BAR_STACKED_100 (0.15 points)
    try:
        # BAR_STACKED_100 has enum value 59
        chart_type_val = int(chart.chart_type)
        if chart_type_val == 59:
            print("PASS: Component 2 -- Chart type is BAR_STACKED_100 (0.15 pts)")
            total_score += 0.15
        else:
            print("FAIL: Component 2 -- Expected BAR_STACKED_100 (59), found chart type {}".format(chart_type_val))
    except Exception as e:
        print("ERROR: Component 2 -- {}".format(e))

    # Component 3: Chart title matches 'Revenue Mix Shift Toward Subscription' (0.15 points)
    try:
        expected_title = "Revenue Mix Shift Toward Subscription"
        if chart.has_title:
            actual_title = chart.chart_title.text_frame.text.strip()
            if actual_title.lower() == expected_title.lower():
                print("PASS: Component 3 -- Chart title matches: '{}'  (0.15 pts)".format(actual_title))
                total_score += 0.15
            else:
                print("FAIL: Component 3 -- Expected title '{}', found '{}'".format(expected_title, actual_title))
        else:
            print("FAIL: Component 3 -- Chart has no title")
    except Exception as e:
        print("ERROR: Component 3 -- {}".format(e))

    # Component 4: Chart data correctness (0.25 points)
    # Expected: 3 series, 4 categories (Q1-Q4), specific percentage values
    # Series 0 (Subscription): [60, 65, 68, 72]
    # Series 1 (Services):     [25, 22, 20, 18]
    # Series 2 (Licensing):    [15, 13, 12, 10]
    try:
        expected_series_values = [
            [60.0, 65.0, 68.0, 72.0],
            [25.0, 22.0, 20.0, 18.0],
            [15.0, 13.0, 12.0, 10.0],
        ]
        expected_categories = ['Q1', 'Q2', 'Q3', 'Q4']

        num_series = len(chart.series)
        data_score = 0.0

        if num_series == 3:
            # Check categories
            try:
                cats = [str(c) for c in chart.plots[0].categories]
            except Exception:
                cats = []

            cats_match = (cats == expected_categories)

            # Check series values with tolerance
            all_values_match = True
            for si in range(3):
                actual_vals = list(chart.series[si].values)
                expected_vals = expected_series_values[si]
                if len(actual_vals) != len(expected_vals):
                    all_values_match = False
                    print("  Series {} length mismatch: expected {}, got {}".format(si, len(expected_vals), len(actual_vals)))
                    break
                for qi in range(len(expected_vals)):
                    if actual_vals[qi] is None or abs(actual_vals[qi] - expected_vals[qi]) > 0.5:
                        all_values_match = False
                        print("  Series {} Q{} mismatch: expected {}, got {}".format(si, qi + 1, expected_vals[qi], actual_vals[qi]))
                        break
                if not all_values_match:
                    break

            if cats_match and all_values_match:
                data_score = 0.25
            elif all_values_match:
                # Values correct but categories might differ slightly
                data_score = 0.20
            elif cats_match:
                data_score = 0.05

            if data_score > 0:
                print("PASS: Component 4 -- Chart data correct (categories_match={}, values_match={}) ({} pts)".format(cats_match, all_values_match, data_score))
            else:
                print("FAIL: Component 4 -- Chart data incorrect (categories_match={}, values_match={})".format(cats_match, all_values_match))
            total_score += data_score
        else:
            print("FAIL: Component 4 -- Expected 3 series, found {}".format(num_series))
    except Exception as e:
        print("ERROR: Component 4 -- {}".format(e))

    # Component 5: Chart series colors (0.25 points)
    # Subscription=#003366, Services=#2196F3, Licensing=#90CAF9
    try:
        expected_colors = ['003366', '2196F3', '90CAF9']
        color_matches = 0
        total_series_for_color = min(len(chart.series), 3)

        for si in range(total_series_for_color):
            try:
                fill = chart.series[si].format.fill
                if fill.type is not None:
                    actual_rgb = str(fill.fore_color.rgb).upper()
                    expected_rgb = expected_colors[si].upper()
                    if actual_rgb == expected_rgb:
                        color_matches += 1
                        print("  Series {} color MATCH: {}".format(si, actual_rgb))
                    else:
                        print("  Series {} color MISMATCH: expected {}, got {}".format(si, expected_rgb, actual_rgb))
                else:
                    print("  Series {} has no solid fill".format(si))
            except Exception as e2:
                print("  Series {} color check error: {}".format(si, e2))

        if color_matches == 3:
            color_score = 0.25
        elif color_matches == 2:
            color_score = 0.15
        elif color_matches == 1:
            color_score = 0.08
        else:
            color_score = 0.0

        if color_score > 0:
            print("PASS: Component 5 -- {} of 3 series colors match ({} pts)".format(color_matches, color_score))
        else:
            print("FAIL: Component 5 -- No series colors match expected values")
        total_score += color_score
    except Exception as e:
        print("ERROR: Component 5 -- {}".format(e))

    final_score = min(total_score, 1.0)
    print("\nScore: {}/1.0".format(total_score))
    print("REWARD: {}".format(final_score))
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = '{}/{}.pptx'.format(WORKDIR, TASK_ID)
if not os.path.exists(file_path):
    print("File not found: {}".format(file_path))
    print("REWARD: 0.0")
else:
    verify_task(file_path)
