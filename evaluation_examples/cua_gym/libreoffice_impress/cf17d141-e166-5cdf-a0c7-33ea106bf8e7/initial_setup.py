"""
Initial Setup: Create Business_Model presentation with 8 slides.
Slide 5 has title 'Revenue Mix Evolution' but NO chart.
Task ID: impress_exec_089
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_089'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_body_slide(prs, title_text, body_lines, layout_idx=1):
    """Add a slide with title and bullet-point body content."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title_text
    # Fill body placeholder
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
    return slide


def add_title_only_slide(prs, title_text, layout_idx=5):
    """Add a blank slide with a manually placed title text box."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x00, 0x33, 0x66)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Business Model Overview"
    slide1.placeholders[1].text = "FY2025 Strategic Revenue Analysis\nPrepared by Finance & Strategy Team"

    # --- Slide 2: Executive Summary ---
    add_title_body_slide(prs, "Executive Summary", [
        "Total annual revenue reached $48.2M, up 23% year-over-year",
        "Subscription revenue now accounts for 66% of total revenue",
        "Customer retention rate improved to 94.2% from 91.5%",
        "Professional services margin expanded by 340 basis points",
        "New enterprise deals contributed $8.7M in incremental ARR",
    ])

    # --- Slide 3: Market Analysis ---
    add_title_body_slide(prs, "Market Analysis", [
        "Total addressable market estimated at $12.4B globally",
        "Cloud-first adoption accelerating across mid-market segment",
        "Key competitors shifting to consumption-based pricing",
        "Regulatory changes in EU creating new compliance opportunities",
        "Partner channel generating 28% of qualified pipeline",
    ])

    # --- Slide 4: Product Strategy ---
    add_title_body_slide(prs, "Product Strategy", [
        "Platform consolidation reducing time-to-value by 40%",
        "AI-powered analytics module launching in Q3 2025",
        "Mobile-first redesign increasing user engagement by 55%",
        "API marketplace enabling third-party integrations",
        "Enterprise tier adding SOC 2 Type II and HIPAA compliance",
    ])

    # --- Slide 5: Revenue Mix Evolution (NO CHART — task target) ---
    slide5 = add_title_only_slide(prs, "Revenue Mix Evolution")
    # Add a subtitle/description text box only, no chart
    txBox = slide5.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Revenue breakdown by segment across fiscal quarters"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # --- Slide 6: Customer Segments ---
    add_title_body_slide(prs, "Customer Segments", [
        "Enterprise (>1000 employees): 42% of revenue, avg deal $285K",
        "Mid-Market (100-999 employees): 35% of revenue, avg deal $48K",
        "SMB (<100 employees): 23% of revenue, avg deal $8.2K",
        "Government & Education vertical growing at 31% CAGR",
        "Healthcare vertical added 47 new accounts in FY2025",
    ])

    # --- Slide 7: Growth Projections ---
    add_title_body_slide(prs, "Growth Projections", [
        "FY2026 revenue target: $62M (29% year-over-year growth)",
        "Subscription mix expected to reach 78% by Q4 FY2026",
        "International expansion targeting APAC and LATAM markets",
        "R&D investment increasing to 22% of revenue",
        "Operating margin projected to improve to 18.5%",
    ])

    # --- Slide 8: Key Takeaways ---
    add_title_body_slide(prs, "Key Takeaways", [
        "Subscription-first strategy driving predictable revenue growth",
        "Platform approach reducing churn and increasing expansion revenue",
        "Strong unit economics: LTV/CAC ratio at 4.8x",
        "Balanced growth across all customer segments",
        "Well-positioned for continued market share gains in FY2026",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
