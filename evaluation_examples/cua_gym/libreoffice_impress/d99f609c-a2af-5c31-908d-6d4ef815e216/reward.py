"""
Reward Script: Sales Enablement Deck (Sales_Playbook.pptx)
Task ID: impress_wf_073
Domain: libreoffice_impress
Scoring:
  C1: File exists on Desktop + 15 slides (0.15)
  C2: Slide 1 title contains 'Sales Playbook' and 'Enterprise Solutions' (0.10)
  C3: Slide 2 has hyperlinks in text (0.10)
  C4: Slide 4 has value proposition canvas with 6 labeled sections (0.10)
  C5: Slide 5 has a competitive battlecard table (0.10)
  C6: Slide 6 has Q&A format with green answer boxes (0.10)
  C7: Slide 7 has flowchart with diamond decision shapes and connectors (0.10)
  C8: Slide 10 has metric value cards (0.10)
  C9: Slide 14 has process flow with arrows and connected shapes (0.10)
  C10: Navy (#1A237E) and gold (#FFB300) colors used (0.05)
  Total: 1.0
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_073'
FILE_PATH = os.path.join(WORKDIR, 'Desktop', 'Sales_Playbook.pptx')


def has_navy_color(prs):
    """Check if navy #1A237E appears in shape fills or text colors across slides."""
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            # Check shape fill
            try:
                if shape.fill.type == 1 and str(shape.fill.fore_color.rgb) == '1A237E':
                    count += 1
            except:
                pass
            # Check text color
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color.type is not None and str(run.font.color.rgb) == '1A237E':
                                count += 1
                                break
                        except:
                            pass
    return count


def has_gold_color(prs):
    """Check if gold #FFB300 appears in shape fills or text colors across slides."""
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if shape.fill.type == 1 and str(shape.fill.fore_color.rgb) == 'FFB300':
                    count += 1
            except:
                pass
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color.type is not None and str(run.font.color.rgb) == 'FFB300':
                                count += 1
                                break
                        except:
                            pass
    return count


def count_hyperlinks_on_slide(slide):
    """Count runs with hyperlink elements on a slide."""
    count = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    rPr = run._r.find(qn('a:rPr'))
                    if rPr is not None:
                        hlinkClick = rPr.find(qn('a:hlinkClick'))
                        if hlinkClick is not None:
                            count += 1
    return count


def get_all_text(slide):
    """Get all text content from a slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    return ' '.join(texts)


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: File on Desktop with 15 slides (0.15 points)
    try:
        num_slides = len(prs.slides)
        if num_slides == 15:
            print(f"PASS: Component 1 -- 15 slides found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 -- Expected 15 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Slide 1 title 'Sales Playbook - Enterprise Solutions' (0.10 points)
    try:
        if num_slides >= 1:
            slide1_text = get_all_text(prs.slides[0]).lower()
            has_playbook = 'sales playbook' in slide1_text
            has_enterprise = 'enterprise solutions' in slide1_text
            if has_playbook and has_enterprise:
                print(f"PASS: Component 2 -- Slide 1 has 'Sales Playbook' and 'Enterprise Solutions' (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 2 -- Slide 1 missing title keywords. playbook={has_playbook}, enterprise={has_enterprise}")
        else:
            print(f"FAIL: Component 2 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Slide 2 has text hyperlinks (0.10 points)
    try:
        if num_slides >= 2:
            hyperlink_count = count_hyperlinks_on_slide(prs.slides[1])
            if hyperlink_count >= 5:
                print(f"PASS: Component 3 -- Slide 2 has {hyperlink_count} hyperlinks (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 3 -- Slide 2 has only {hyperlink_count} hyperlinks, expected >=5")
        else:
            print(f"FAIL: Component 3 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Slide 4 value proposition canvas with 6 labeled sections (0.10 points)
    try:
        if num_slides >= 4:
            slide4 = prs.slides[3]
            slide4_text = get_all_text(slide4).lower()
            # Check for the 6 sections of value proposition canvas
            sections = ['customer jobs', 'pains', 'gains', 'products', 'pain relievers', 'gain creators']
            found_sections = [s for s in sections if s in slide4_text]
            if len(found_sections) >= 5:
                print(f"PASS: Component 4 -- Slide 4 has {len(found_sections)}/6 VPC sections: {found_sections} (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 4 -- Slide 4 only has {len(found_sections)}/6 VPC sections: {found_sections}")
        else:
            print(f"FAIL: Component 4 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Slide 5 has competitive battlecard table (0.10 points)
    try:
        if num_slides >= 5:
            slide5 = prs.slides[4]
            has_table = False
            table_rows = 0
            table_cols = 0
            for shape in slide5.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    has_table = True
                    table_rows = len(shape.table.rows)
                    table_cols = len(shape.table.columns)
                    break
            slide5_text = get_all_text(slide5).lower()
            has_battlecard_title = 'battlecard' in slide5_text or 'competitive' in slide5_text
            if has_table and has_battlecard_title and table_rows >= 3 and table_cols >= 3:
                print(f"PASS: Component 5 -- Slide 5 has battlecard table ({table_rows}x{table_cols}) (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 5 -- table={has_table}, title={has_battlecard_title}, rows={table_rows}, cols={table_cols}")
        else:
            print(f"FAIL: Component 5 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # Component 6: Slide 6 Q&A with green answer boxes (0.10 points)
    try:
        if num_slides >= 6:
            slide6 = prs.slides[5]
            slide6_text = get_all_text(slide6)
            # Check for Q: and A: pattern
            has_questions = slide6_text.count('Q:') >= 2 or slide6_text.count('Q.') >= 2
            has_answers = slide6_text.count('A:') >= 2 or slide6_text.count('A.') >= 2
            # Check for green-ish filled shapes (answer boxes)
            green_boxes = 0
            for shape in slide6.shapes:
                try:
                    if shape.fill.type == 1:  # SOLID
                        rgb = str(shape.fill.fore_color.rgb)
                        # Check for green shades: E8F5E9, 4CAF50, 2E7D32, 00C853, etc.
                        r = int(rgb[0:2], 16)
                        g = int(rgb[2:4], 16)
                        b = int(rgb[4:6], 16)
                        if g > r and g > b:  # green dominant
                            green_boxes += 1
                except:
                    pass
            if has_questions and has_answers and green_boxes >= 2:
                print(f"PASS: Component 6 -- Slide 6 has Q&A format with {green_boxes} green boxes (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 6 -- questions={has_questions}, answers={has_answers}, green_boxes={green_boxes}")
        else:
            print(f"FAIL: Component 6 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    # Component 7: Slide 7 flowchart with diamond decision shapes and connectors (0.10 points)
    try:
        if num_slides >= 7:
            slide7 = prs.slides[6]
            diamonds = 0
            connectors = 0
            rounded_rects = 0
            for shape in slide7.shapes:
                if 'Diamond' in shape.name:
                    diamonds += 1
                if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM or shape.shape_type == 9:  # LINE type
                    connectors += 1
                if 'Connector' in shape.name:
                    connectors += 1
                if 'Rounded Rectangle' in shape.name:
                    rounded_rects += 1
            # Deduplicate connector count (some may match both type and name)
            slide7_text = get_all_text(slide7).lower()
            has_pricing = 'pricing' in slide7_text or 'decision' in slide7_text
            if diamonds >= 1 and connectors >= 2 and has_pricing:
                print(f"PASS: Component 7 -- Slide 7 flowchart: {diamonds} diamonds, {connectors} connectors (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 7 -- diamonds={diamonds}, connectors={connectors}, pricing_text={has_pricing}")
        else:
            print(f"FAIL: Component 7 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 7 -- {e}")

    # Component 8: Slide 10 has metric value cards (0.10 points)
    try:
        if num_slides >= 10:
            slide10 = prs.slides[9]
            slide10_text = get_all_text(slide10)
            # Check for metric values like percentages, multipliers, dollar amounts
            has_percent = '%' in slide10_text
            has_multiplier = 'x' in slide10_text.lower() and ('roi' in slide10_text.lower() or '3.2' in slide10_text)
            has_dollar = '$' in slide10_text
            # Check for case study context
            has_case_study = 'case study' in slide10_text.lower() or 'acme' in slide10_text.lower()
            metric_count = sum([has_percent, has_multiplier, has_dollar])
            if metric_count >= 2 and has_case_study:
                print(f"PASS: Component 8 -- Slide 10 has metric cards: percent={has_percent}, multiplier={has_multiplier}, dollar={has_dollar} (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 8 -- metrics={metric_count}, case_study={has_case_study}")
        else:
            print(f"FAIL: Component 8 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 8 -- {e}")

    # Component 9: Slide 14 has process flow with arrows and connected shapes (0.10 points)
    try:
        if num_slides >= 14:
            slide14 = prs.slides[13]
            arrows = 0
            flow_shapes = 0
            for shape in slide14.shapes:
                if 'Arrow' in shape.name:
                    arrows += 1
                if 'Rounded Rectangle' in shape.name:
                    flow_shapes += 1
            slide14_text = get_all_text(slide14).lower()
            has_handoff = 'handoff' in slide14_text or 'process' in slide14_text
            if arrows >= 3 and flow_shapes >= 4 and has_handoff:
                print(f"PASS: Component 9 -- Slide 14 process flow: {arrows} arrows, {flow_shapes} flow shapes (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 9 -- arrows={arrows}, flow_shapes={flow_shapes}, handoff={has_handoff}")
        else:
            print(f"FAIL: Component 9 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 9 -- {e}")

    # Component 10: Navy (#1A237E) and gold (#FFB300) colors used throughout (0.05 points)
    try:
        navy_count = has_navy_color(prs)
        gold_count = has_gold_color(prs)
        if navy_count >= 5 and gold_count >= 2:
            print(f"PASS: Component 10 -- Colors: navy={navy_count} instances, gold={gold_count} instances (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 10 -- navy={navy_count}, gold={gold_count} (need navy>=5, gold>=2)")
    except Exception as e:
        print(f"ERROR: Component 10 -- {e}")

    final_score = min(total_score, 1.0)
    final_score = round(final_score, 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for unsaved GUI edits
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(FILE_PATH)
