"""
Initial Setup: 5-slide product demo deck with images in various non-top positions and unformatted body text
Task ID: osworld_impress_image_top_underline_text_009
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import io
from PIL import Image, ImageDraw

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_image_top_underline_text_009'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def make_image_bytes(color: tuple, width: int = 200, height: int = 150, label: str = "") -> bytes:
    """Generate a simple colored PNG image as bytes."""
    img = Image.new("RGB", (width, height), color=color)
    draw = ImageDraw.Draw(img)
    draw.rectangle([5, 5, width - 6, height - 6], outline=(255, 255, 255), width=3)
    if label:
        draw.text((10, 10), label, fill=(255, 255, 255))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide data: title, body text, image color, image position (NON-TOP: middle or bottom)
    slides_data = [
        {
            "title": "NovaTech X1 Pro — Product Overview",
            "body_lines": [
                "The NovaTech X1 Pro redefines portable computing with its 13-inch OLED display.",
                "Weighing just 1.1 kg, it delivers all-day battery life up to 18 hours.",
                "Available in Midnight Silver and Cosmic Blue finishes.",
            ],
            "img_color": (41, 98, 255),
            "img_label": "X1 Pro",
            # Image placed in middle area
            "img_left": Inches(6.5),
            "img_top": Inches(3.0),
            "img_width": Inches(3.0),
            "img_height": Inches(2.2),
        },
        {
            "title": "Performance & Processing Power",
            "body_lines": [
                "Powered by the latest 12-core NovaTech Helix processor clocked at 3.8 GHz.",
                "Integrated NovaTech GPU handles 4K video editing with ease.",
                "32 GB LPDDR5 RAM ensures seamless multitasking across demanding workflows.",
            ],
            "img_color": (220, 53, 69),
            "img_label": "Performance",
            # Image placed in bottom area
            "img_left": Inches(0.5),
            "img_top": Inches(5.0),
            "img_width": Inches(3.5),
            "img_height": Inches(2.0),
        },
        {
            "title": "Connectivity & Expansion",
            "body_lines": [
                "Features USB-C Thunderbolt 4 (x2), USB-A 3.2, and HDMI 2.1 ports.",
                "Wi-Fi 7 and Bluetooth 5.4 ensure fast, reliable wireless connectivity.",
                "Nano-SIM slot supports optional 5G connectivity for mobile professionals.",
            ],
            "img_color": (40, 167, 69),
            "img_label": "Connectivity",
            # Image in middle-right area
            "img_left": Inches(6.8),
            "img_top": Inches(2.8),
            "img_width": Inches(2.8),
            "img_height": Inches(2.1),
        },
        {
            "title": "Display & Audio Experience",
            "body_lines": [
                "13.3-inch OLED panel with 2560x1600 resolution and 120Hz adaptive refresh rate.",
                "Supports 100% DCI-P3 color gamut with Delta-E < 1 factory calibration.",
                "Quad-speaker Dolby Atmos system delivers immersive spatial audio.",
            ],
            "img_color": (255, 153, 0),
            "img_label": "Display",
            # Image in bottom-center area
            "img_left": Inches(3.5),
            "img_top": Inches(5.2),
            "img_width": Inches(3.0),
            "img_height": Inches(2.0),
        },
        {
            "title": "Pricing & Availability",
            "body_lines": [
                "Base model starts at $1,299 with 16 GB RAM and 512 GB NVMe SSD.",
                "Pro configuration ($1,799) includes 32 GB RAM and 1 TB SSD.",
                "Available at NovaTech.com and authorized retailers from April 15, 2025.",
            ],
            "img_color": (111, 66, 193),
            "img_label": "Pricing",
            # Image in center area
            "img_left": Inches(6.0),
            "img_top": Inches(3.5),
            "img_width": Inches(3.2),
            "img_height": Inches(2.2),
        },
    ]

    for slide_data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

        # Title text box
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9.0), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = False
        p_title = tf_title.paragraphs[0]
        run_title = p_title.add_run()
        run_title.text = slide_data["title"]
        run_title.font.bold = True
        run_title.font.size = Pt(24)
        run_title.font.color.rgb = RGBColor(0x1F, 0x35, 0x64)

        # Body text box (below title area, NOT underlined in initial state)
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.5), Inches(3.5))
        tf_body = body_box.text_frame
        tf_body.word_wrap = True
        for i, line in enumerate(slide_data["body_lines"]):
            if i == 0:
                p = tf_body.paragraphs[0]
            else:
                p = tf_body.add_paragraph()
            p.space_after = Pt(6)
            run = p.add_run()
            run.text = line
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            # NOTE: No underline in initial state

        # Image in non-top position
        img_bytes = make_image_bytes(
            slide_data["img_color"],
            label=slide_data["img_label"]
        )
        slide.shapes.add_picture(
            img_bytes,
            slide_data["img_left"],
            slide_data["img_top"],
            slide_data["img_width"],
            slide_data["img_height"],
        )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
