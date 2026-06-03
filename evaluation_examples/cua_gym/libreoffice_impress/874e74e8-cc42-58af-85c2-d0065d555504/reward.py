"""
Reward Script: Move images to top area and underline all body text
Task ID: osworld_impress_image_top_underline_text_009
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.5 pts): All 5 images moved to top area of their slides
                            (image.top < 25% of slide height)
  - Component 2 (0.5 pts): All text runs in all 5 body textboxes are underlined
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_image_top_underline_text_009'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0

    Task requirements:
    1. Move all images on all slides to the top area (top < 25% of slide height)
    2. Underline all text in all body textboxes (TextBox 3) on all 5 slides
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print("CRITICAL: Cannot load file %s: %s" % (file_path, e))
        print("REWARD: 0.0")
        return 0.0

    slide_height = prs.slide_height  # 6858000 EMU for 7.5 inch slides
    # "Top area" = top 25% of slide height
    top_area_threshold = int(slide_height * 0.25)  # 1714500 EMU ~ 1.875 inches
    num_slides = len(prs.slides)

    print("File: %s" % file_path)
    print("Slide height: %d EMU (%.2f inches)" % (slide_height, slide_height / 914400))
    print("Top area threshold: %d EMU (%.2f inches)" % (top_area_threshold, top_area_threshold / 914400))
    print("Number of slides: %d" % num_slides)
    print()

    # Component 1: All images moved to top area of their slides (0.5 points)
    # Each image top must be < 25% of slide height
    # In the initial state, images are at 2.8 to 5.2 inches (well below threshold)
    # In the golden state, images are at 0.1 inch (well within threshold)
    try:
        images_in_top = 0
        images_total = 0
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    images_total += 1
                    img_top = shape.top if shape.top is not None else 0
                    in_top = img_top < top_area_threshold
                    print("Slide %d image top: %d EMU (%.3f in) -> %s" % (
                        i + 1, img_top, img_top / 914400, "TOP AREA" if in_top else "NOT in top area"
                    ))
                    if in_top:
                        images_in_top += 1

        print("Images in top area: %d / %d" % (images_in_top, images_total))
        if images_total == 5 and images_in_top == 5:
            print("PASS: Component 1 — All 5 images moved to top area (0.5 pts)")
            total_score += 0.5
        elif images_in_top > 0:
            # Partial credit not applicable in this two-component design;
            # all or nothing per component. But we still log.
            print("FAIL: Component 1 — Only %d/5 images in top area" % images_in_top)
        else:
            print("FAIL: Component 1 — No images in top area (%d/%d)" % (images_in_top, images_total))
    except Exception as e:
        print("ERROR: Component 1 — %s" % e)

    print()

    # Component 2: All text runs in body textboxes (TextBox 3) are underlined (0.5 points)
    # In initial state: all runs have underline=None (not underlined)
    # In golden state: all runs have underline=True
    try:
        total_runs = 0
        underlined_runs = 0
        body_textboxes_found = 0

        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                # Body textboxes are named "TextBox 3" per the task structure
                if shape.has_text_frame and 'TextBox 3' in shape.name:
                    body_textboxes_found += 1
                    for k, para in enumerate(shape.text_frame.paragraphs):
                        for run in para.runs:
                            run_text = (run.text or '').strip()
                            if not run_text:
                                continue  # skip empty runs
                            total_runs += 1
                            is_underlined = run.font.underline is True
                            print("Slide %d TextBox3 Para %d: text=%r... underline=%s -> %s" % (
                                i + 1, k, run_text[:20], run.font.underline, "UNDERLINED" if is_underlined else "NOT underlined"
                            ))
                            if is_underlined:
                                underlined_runs += 1

        print("Body textboxes found: %d" % body_textboxes_found)
        print("Underlined runs: %d / %d" % (underlined_runs, total_runs))
        if body_textboxes_found == 5 and total_runs > 0 and underlined_runs == total_runs:
            print("PASS: Component 2 — All %d text runs in body textboxes are underlined (0.5 pts)" % total_runs)
            total_score += 0.5
        else:
            print("FAIL: Component 2 — Only %d/%d runs are underlined across %d body textboxes" % (
                underlined_runs, total_runs, body_textboxes_found
            ))
    except Exception as e:
        print("ERROR: Component 2 — %s" % e)

    print()
    final_score = min(total_score, 1.0)
    print("Score: %.1f/1.0" % total_score)
    print("REWARD: %.1f" % final_score)
    return final_score


# Default: test against canonical artifact path on the VM
file_path = '%s/%s.pptx' % (WORKDIR, TASK_ID)
if not os.path.exists(file_path):
    print("File not found: %s" % file_path)
    print("REWARD: 0.0")
else:
    verify_task(file_path)
