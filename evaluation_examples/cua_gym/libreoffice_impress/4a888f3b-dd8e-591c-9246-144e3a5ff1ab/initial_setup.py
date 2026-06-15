"""
Initial Setup: Create a photo album presentation with 6 slides
Task ID: impstruct_044
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
OUTPUT = f'{WORKDIR}/photo_album.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_colored_rect(slide, left, top, width, height, r, g, b):
    """Add a colored rectangle shape as a stand-in for a photo."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()
    return shape


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    fill1 = slide1.background.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    txb = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Coastal Adventures 2025"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    p2 = tf.add_paragraph()
    p2.text = "A Photo Journey Along the Pacific Northwest"
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.runs[0]
    r2.font.name = "Arial"
    r2.font.size = Pt(20)
    r2.font.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)

    # --- Slide 2: Sunrise at Cannon Beach ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    txb2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf2 = txb2.text_frame
    p = tf2.paragraphs[0]
    p.text = "Sunrise at Cannon Beach"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xE7, 0x4C, 0x3C)

    # Large "photo" rectangle
    add_colored_rect(slide2, Inches(0.5), Inches(1.5), Inches(6), Inches(4.5),
                     0xF3, 0x9C, 0x12)

    # Caption
    txb2c = slide2.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(6), Inches(0.8))
    tf2c = txb2c.text_frame
    p = tf2c.paragraphs[0]
    p.text = "Haystack Rock silhouetted against the morning sky, captured at 6:47 AM"
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(14)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)

    # --- Slide 3: Tide Pool Discoveries ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txb3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf3 = txb3.text_frame
    p = tf3.paragraphs[0]
    p.text = "Tide Pool Discoveries"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x27, 0xAE, 0x60)

    # Grid of 4 "photos"
    positions = [
        (Inches(0.5), Inches(1.5), Inches(4), Inches(2.5)),
        (Inches(5), Inches(1.5), Inches(4), Inches(2.5)),
        (Inches(0.5), Inches(4.3), Inches(4), Inches(2.5)),
        (Inches(5), Inches(4.3), Inches(4), Inches(2.5)),
    ]
    colors = [
        (0x1A, 0xBC, 0x9C), (0x29, 0x80, 0xB9),
        (0x8E, 0x44, 0xAD), (0x2E, 0xCC, 0x71),
    ]
    labels = ["Sea Anemone", "Starfish Colony", "Hermit Crab", "Kelp Forest"]
    for (l, t, w, h), (cr, cg, cb), label in zip(positions, colors, labels):
        add_colored_rect(slide3, l, t, w, h, cr, cg, cb)
        ltxb = slide3.shapes.add_textbox(l, t + h + Emu(36000), w, Inches(0.4))
        ltf = ltxb.text_frame
        lp = ltf.paragraphs[0]
        lp.text = label
        lp.alignment = PP_ALIGN.CENTER
        lr = lp.runs[0]
        lr.font.name = "Arial"
        lr.font.size = Pt(12)
        lr.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # --- Slide 4: Hiking the Oregon Coast Trail ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    fill4 = slide4.background.fill
    fill4.solid()
    fill4.fore_color.rgb = RGBColor(0xEC, 0xF0, 0xF1)

    txb4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf4 = txb4.text_frame
    p = tf4.paragraphs[0]
    p.text = "Hiking the Oregon Coast Trail"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)

    add_colored_rect(slide4, Inches(1.5), Inches(1.5), Inches(7), Inches(4),
                     0x16, 0xA0, 0x85)

    txb4d = slide4.shapes.add_textbox(Inches(1), Inches(5.8), Inches(8), Inches(1.2))
    tf4d = txb4d.text_frame
    tf4d.word_wrap = True
    p = tf4d.paragraphs[0]
    p.text = ("Day 3 of our trek: 12.4 miles from Ecola State Park to Hug Point. "
              "The coastal fog cleared around noon revealing stunning cliff views.")
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # --- Slide 5: Wildlife Encounters ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    txb5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf5 = txb5.text_frame
    p = tf5.paragraphs[0]
    p.text = "Wildlife Encounters"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xD3, 0x54, 0x00)

    # Two images side by side
    add_colored_rect(slide5, Inches(0.5), Inches(1.5), Inches(4.2), Inches(3.5),
                     0x5D, 0x6D, 0x7E)
    add_colored_rect(slide5, Inches(5.3), Inches(1.5), Inches(4.2), Inches(3.5),
                     0x85, 0xC1, 0xE9)

    txb5a = slide5.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(4.2), Inches(0.5))
    tf5a = txb5a.text_frame
    p = tf5a.paragraphs[0]
    p.text = "Gray Whale Breach"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(14)
    run.font.bold = True

    txb5b = slide5.shapes.add_textbox(Inches(5.3), Inches(5.2), Inches(4.2), Inches(0.5))
    tf5b = txb5b.text_frame
    p = tf5b.paragraphs[0]
    p.text = "Bald Eagle in Flight"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(14)
    run.font.bold = True

    txb5c = slide5.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(1))
    tf5c = txb5c.text_frame
    tf5c.word_wrap = True
    p = tf5c.paragraphs[0]
    p.text = ("Spotted three gray whales migrating north and a nesting pair of bald eagles "
              "near Cape Meares lighthouse.")
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(13)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)

    # --- Slide 6: Thank You ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    fill6 = slide6.background.fill
    fill6.solid()
    fill6.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    txb6 = slide6.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf6 = txb6.text_frame
    tf6.word_wrap = True
    p = tf6.paragraphs[0]
    p.text = "Thank You for Viewing"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    p2 = tf6.add_paragraph()
    p2.text = "Photography by Elena Vasquez | June 2025"
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.runs[0]
    r2.font.name = "Arial"
    r2.font.size = Pt(18)
    r2.font.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)

    p3 = tf6.add_paragraph()
    p3.text = "www.elenavasquez-photography.com"
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.runs[0]
    r3.font.name = "Arial"
    r3.font.size = Pt(16)
    r3.font.color.rgb = RGBColor(0x3A, 0x97, 0xD4)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open file in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
