"""
Initial Setup: Art History Lecture Presentation
Task ID: impress_wf_067
Domain: libreoffice_impress

Creates a blank presentation and opens it in LibreOffice Impress.
The agent's task is to create the full Art_History.pptx from scratch.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_067'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    # Create a blank presentation (default single blank slide)
    prs = Presentation()
    # Add a single blank slide so Impress opens with something
    blank_layout = prs.slide_layouts[6]  # Title Only or Blank
    prs.slides.add_slide(blank_layout)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
