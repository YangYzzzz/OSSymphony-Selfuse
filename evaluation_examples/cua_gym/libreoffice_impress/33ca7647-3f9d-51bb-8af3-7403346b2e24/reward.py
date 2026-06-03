"""
Reward Script: Art History Lecture Presentation
Task ID: impress_wf_067
Domain: libreoffice_impress
Scoring:
  Component 1 (0.15): File exists at Desktop/Art_History.pptx with 10 slides
  Component 2 (0.15): Slide 1 title is 'Art Through the Ages'
  Component 3 (0.25): Slides 2-8 each have period title, 3 rectangle placeholders, and callout shape
  Component 4 (0.15): Slide 9 has arrow + 7 colored rectangle blocks (timeline)
  Component 5 (0.15): Slide 10 has 'Discussion Questions'
  Component 6 (0.15): Background #FFF3E0 and text color #4E342E across slides
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_067'
FILE_PATH = os.path.join(WORKDIR, 'Desktop', 'Art_History.pptx')

# Expected art periods for slides 2-8
PERIODS = [
    'Renaissance',
    'Baroque',
    'Impressionism',
    'Post-Impressionism',
    'Modernism',
    'Abstract Expressionism',
    'Contemporary',
]


def get_all_text(slide):
    """Recursively extract all text from a slide, including grouped shapes."""
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text:
            texts.append(shape.text)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                if hasattr(sub, 'text') and sub.text:
                    texts.append(sub.text)
    return texts


def get_slide_bg_color(slide):
    """Get slide background RGB color, handling inherited fills."""
    try:
        fill = slide.background.fill
        if fill.type == 1:
            return str(fill.fore_color.rgb)
        elif fill.type == 5:
            master_fill = slide.slide_layout.slide_master.background.fill
            if master_fill.type == 1:
                return str(master_fill.fore_color.rgb)
    except Exception:
        pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: File has exactly 10 slides (0.15 points)
    try:
        if num_slides == 10:
            print(f"PASS: Component 1 - File has 10 slides (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 - Expected 10 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Slide 1 title is 'Art Through the Ages' (0.15 points)
    try:
        if num_slides >= 1:
            slide1_texts = get_all_text(prs.slides[0])
            slide1_all = ' '.join(slide1_texts).lower()
            if 'art through the ages' in slide1_all:
                print(f"PASS: Component 2 - Slide 1 contains 'Art Through the Ages' (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 2 - 'Art Through the Ages' not found in slide 1. Found texts: {slide1_texts[:3]}")
        else:
            print(f"FAIL: Component 2 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Slides 2-8 each have period title, 3 rectangle placeholders, and callout (0.25 points)
    # Each of the 7 period slides is worth ~0.036 points; we need all 7 to get full credit
    try:
        if num_slides >= 8:
            period_slides_ok = 0
            for i in range(7):
                slide = prs.slides[i + 1]  # slides 2-8 = indices 1-7
                period = PERIODS[i]

                # Check period name appears in the slide text
                slide_texts = get_all_text(slide)
                slide_all_text = ' '.join(slide_texts).lower()
                has_period_name = period.lower() in slide_all_text

                # Count rectangle auto shapes (image placeholders) - exclude rounded rectangles
                rect_count = 0
                has_callout = False
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        name_lower = shape.name.lower()
                        if 'rounded' in name_lower:
                            # This is the callout box
                            has_callout = True
                        elif 'rectangle' in name_lower:
                            rect_count += 1

                slide_ok = has_period_name and rect_count >= 3 and has_callout
                if slide_ok:
                    period_slides_ok += 1
                    print(f"  Slide {i+2} ({period}): OK (period_name={has_period_name}, rects={rect_count}, callout={has_callout})")
                else:
                    print(f"  Slide {i+2} ({period}): INCOMPLETE (period_name={has_period_name}, rects={rect_count}, callout={has_callout})")

            if period_slides_ok == 7:
                print(f"PASS: Component 3 - All 7 period slides have required elements (0.25 pts)")
                total_score += 0.25
            elif period_slides_ok >= 4:
                partial = round(0.25 * period_slides_ok / 7, 2)
                print(f"PARTIAL: Component 3 - {period_slides_ok}/7 period slides OK ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 3 - Only {period_slides_ok}/7 period slides have required elements")
        else:
            print(f"FAIL: Component 3 - Not enough slides (need >= 8)")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Slide 9 has timeline with arrow + 7 colored rectangle blocks (0.15 points)
    try:
        if num_slides >= 9:
            slide9 = prs.slides[8]
            has_arrow = False
            colored_rects = 0

            for shape in slide9.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    name_lower = shape.name.lower()
                    if 'arrow' in name_lower:
                        has_arrow = True
                    elif 'rectangle' in name_lower:
                        # Check it has a solid fill (colored block)
                        try:
                            if shape.fill.type == 1:
                                colored_rects += 1
                        except Exception:
                            colored_rects += 1  # count it anyway if fill check fails

            if has_arrow and colored_rects >= 7:
                print(f"PASS: Component 4 - Slide 9 timeline: arrow={has_arrow}, colored_rects={colored_rects} (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 4 - Slide 9 timeline: arrow={has_arrow}, colored_rects={colored_rects} (need arrow + 7 rects)")
        else:
            print(f"FAIL: Component 4 - Not enough slides (need >= 9)")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Slide 10 has 'Discussion Questions' (0.15 points)
    try:
        if num_slides >= 10:
            slide10_texts = get_all_text(prs.slides[9])
            slide10_all = ' '.join(slide10_texts).lower()
            if 'discussion questions' in slide10_all:
                print(f"PASS: Component 5 - Slide 10 contains 'Discussion Questions' (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 5 - 'Discussion Questions' not found in slide 10. Found: {slide10_texts[:2]}")
        else:
            print(f"FAIL: Component 5 - Not enough slides (need >= 10)")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    # Component 6: Background #FFF3E0 and text color #4E342E across slides (0.15 points)
    try:
        bg_ok_count = 0
        text_color_ok_count = 0
        slides_checked = min(num_slides, 10)

        for i in range(slides_checked):
            slide = prs.slides[i]

            # Check background
            bg_color = get_slide_bg_color(slide)
            if bg_color and bg_color.upper() == 'FFF3E0':
                bg_ok_count += 1

            # Check at least one text shape has #4E342E color
            found_brown = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                if run.font.color.type is not None:
                                    rgb_str = str(run.font.color.rgb).upper()
                                    if rgb_str == '4E342E':
                                        found_brown = True
                                        break
                            except Exception:
                                pass
                        if found_brown:
                            break
                if found_brown:
                    break
            if found_brown:
                text_color_ok_count += 1

        # Need at least 8 out of 10 slides with correct bg and text color
        bg_pass = bg_ok_count >= 8
        text_pass = text_color_ok_count >= 8

        if bg_pass and text_pass:
            print(f"PASS: Component 6 - BG #FFF3E0 on {bg_ok_count}/{slides_checked} slides, text #4E342E on {text_color_ok_count}/{slides_checked} slides (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 6 - BG #FFF3E0 on {bg_ok_count}/{slides_checked} slides (need 8+), text #4E342E on {text_color_ok_count}/{slides_checked} slides (need 8+)")
    except Exception as e:
        print(f"ERROR: Component 6 - {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for unsaved GUI state
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(FILE_PATH)
