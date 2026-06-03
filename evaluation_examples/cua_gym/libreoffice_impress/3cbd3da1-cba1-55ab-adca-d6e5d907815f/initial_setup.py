"""
Initial Setup: Create a 10-slide presentation where slide 5 has 6 bullet points.
Task ID: impress_ndo_089
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_089'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullets, layout_idx=1):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.paragraphs[0].text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(prs,
        "Nexus Analytics - Q4 2025 Strategy Review",
        "Prepared by the Product & Engineering Team\nDecember 2025"
    )

    # Slide 2: Agenda
    add_content_slide(prs, "Agenda", [
        "Market landscape and competitive analysis",
        "Product roadmap highlights for Q4",
        "Engineering milestones and technical debt review",
        "Revenue targets and growth projections",
        "Key features of the new platform release",
        "Team expansion and resource allocation",
        "Risk assessment and mitigation strategies",
        "Open discussion and next steps",
    ])

    # Slide 3: Market Overview
    add_content_slide(prs, "Market Overview", [
        "Global SaaS market projected to reach $819B by 2030",
        "Our segment grew 23% year-over-year in enterprise analytics",
        "Three new competitors entered the mid-market space in Q3",
        "Customer retention rate improved to 94.7% from 91.2%",
        "Average deal size increased by 18% compared to Q3",
    ])

    # Slide 4: Product Roadmap
    add_content_slide(prs, "Product Roadmap - Q4 Priorities", [
        "Launch real-time collaborative dashboards by November 15",
        "Complete migration to microservices architecture",
        "Integrate natural language query engine (beta)",
        "Deploy enhanced data connector framework for 40+ sources",
        "Release mobile application v2.0 with offline capabilities",
    ])

    # Slide 5: Key Features (the target slide with 6 bullet points)
    add_content_slide(prs, "Key Features", [
        "Predictive analytics engine with 95% accuracy on trend forecasting",
        "Real-time data synchronization across all connected platforms",
        "Custom dashboard builder with drag-and-drop widget library",
        "Role-based access control with granular permission management",
        "Automated report generation with scheduling and distribution",
        "Multi-language support covering 28 languages with RTL compatibility",
    ])

    # Slide 6: Engineering Milestones
    add_content_slide(prs, "Engineering Milestones", [
        "Reduced API response time by 40% through query optimization",
        "Achieved 99.97% uptime across all production services",
        "Migrated 78% of monolith services to containerized microservices",
        "Implemented end-to-end encryption for data at rest and in transit",
        "Completed SOC 2 Type II compliance audit successfully",
    ])

    # Slide 7: Revenue Projections
    add_content_slide(prs, "Revenue Projections", [
        "Q4 revenue target: $12.8M (up from $10.4M in Q3)",
        "New enterprise contracts pipeline valued at $4.2M",
        "Expansion revenue from existing accounts: $2.1M projected",
        "Professional services backlog at $1.8M",
        "Annual recurring revenue expected to cross $45M by year end",
    ])

    # Slide 8: Team Expansion
    add_content_slide(prs, "Team Expansion Plan", [
        "Hiring 15 senior engineers for the platform team",
        "Opening new development center in Austin, TX",
        "Establishing dedicated customer success team of 8 specialists",
        "Recruiting VP of Data Science to lead ML initiatives",
        "Planned headcount growth from 142 to 178 by end of Q1 2026",
    ])

    # Slide 9: Risk Assessment
    add_content_slide(prs, "Risk Assessment", [
        "Supply chain delays may impact hardware procurement timeline",
        "Regulatory changes in EU data privacy require compliance updates",
        "Key person dependency in the core infrastructure team",
        "Potential talent acquisition challenges in competitive job market",
        "Third-party API deprecation risk for legacy integrations",
    ])

    # Slide 10: Next Steps
    add_content_slide(prs, "Next Steps & Action Items", [
        "Finalize Q4 sprint planning by October 10",
        "Submit SOC 2 Type II renewal documentation",
        "Complete customer advisory board recruiting",
        "Present updated pricing model to executive leadership",
        "Schedule cross-functional alignment workshop for November",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
