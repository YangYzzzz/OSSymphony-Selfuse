"""
Reward Script: Split slide 5 into two slides - keep first 3 bullets on slide 5,
move last 3 bullets to new slide 6 with '(continued)' title prefix.
Task ID: impress_ndo_089
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): Total slide count is 11 (was 10)
  Component 2 (0.25): Slide 5 has title 'Key Features' with first 3 bullets only
  Component 3 (0.30): Slide 6 has title 'Key Features (continued)' with last 3 bullets
  Component 4 (0.25): Original slides 6-10 shifted to positions 7-11 correctly
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_089'

# Expected first 3 bullets (stay on slide 5)
FIRST_3_BULLETS = [
    'Predictive analytics engine with 95% accuracy on trend forecasting',
    'Real-time data synchronization across all connected platforms',
    'Custom dashboard builder with drag-and-drop widget library',
]

# Expected last 3 bullets (move to new slide 6)
LAST_3_BULLETS = [
    'Role-based access control with granular permission management',
    'Automated report generation with scheduling and distribution',
    'Multi-language support covering 28 languages with RTL compatibility',
]

# Expected titles for slides 7-11 (originally slides 6-10, shifted by 1)
SHIFTED_TITLES = [
    'Engineering Milestones',
    'Revenue Projections',
    'Team Expansion Plan',
    'Risk Assessment',
    'Next Steps & Action Items',
]


def get_slide_title(slide):
    """Extract the title text from a slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    return text
    return ""


def get_slide_bullets(slide):
    """Extract non-title text paragraphs from a slide (skips the first non-empty text as title)."""
    all_texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    all_texts.append(text)
    # First non-empty text is the title; rest are bullets
    return all_texts[1:] if len(all_texts) > 1 else []


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Total slide count is 11 (0.20 points)
    # Initial has 10 slides; golden has 11 (one new slide inserted)
    try:
        if num_slides == 11:
            print(f"PASS: Component 1 — Slide count is 11 (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 — Expected 11 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 5 has title 'Key Features' and first 3 bullets only (0.25 points)
    try:
        if num_slides >= 5:
            slide5 = prs.slides[4]  # 0-indexed
            title5 = get_slide_title(slide5)
            bullets5 = get_slide_bullets(slide5)

            title_ok = title5 == 'Key Features'
            bullets_ok = (len(bullets5) == 3 and
                          all(b == expected for b, expected in zip(bullets5, FIRST_3_BULLETS)))

            if title_ok and bullets_ok:
                print(f"PASS: Component 2 — Slide 5 title 'Key Features' with correct 3 bullets (0.25 pts)")
                total_score += 0.25
            elif title_ok and not bullets_ok:
                # Partial: title correct but bullets wrong
                print(f"FAIL: Component 2 — Title correct but bullets wrong. Found {len(bullets5)} bullets: {bullets5}")
            elif not title_ok and bullets_ok:
                print(f"FAIL: Component 2 — Bullets correct but title wrong: '{title5}'")
            else:
                print(f"FAIL: Component 2 — Title: '{title5}', Bullets ({len(bullets5)}): {bullets5}")
        else:
            print(f"FAIL: Component 2 — Not enough slides ({num_slides})")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 6 has title 'Key Features (continued)' and last 3 bullets (0.30 points)
    try:
        if num_slides >= 6:
            slide6 = prs.slides[5]  # 0-indexed
            title6 = get_slide_title(slide6)
            bullets6 = get_slide_bullets(slide6)

            title_ok = title6 == 'Key Features (continued)'
            bullets_ok = (len(bullets6) == 3 and
                          all(b == expected for b, expected in zip(bullets6, LAST_3_BULLETS)))

            if title_ok and bullets_ok:
                print(f"PASS: Component 3 — Slide 6 title 'Key Features (continued)' with correct 3 bullets (0.30 pts)")
                total_score += 0.30
            elif title_ok and not bullets_ok:
                print(f"FAIL: Component 3 — Title correct but bullets wrong. Found {len(bullets6)} bullets: {bullets6}")
            elif not title_ok and bullets_ok:
                print(f"FAIL: Component 3 — Bullets correct but title wrong: '{title6}'")
            else:
                print(f"FAIL: Component 3 — Title: '{title6}', Bullets ({len(bullets6)}): {bullets6}")
        else:
            print(f"FAIL: Component 3 — Not enough slides ({num_slides})")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Original slides 6-10 shifted to 7-11 with correct titles (0.25 points)
    try:
        if num_slides >= 11:
            matched = 0
            for i, expected_title in enumerate(SHIFTED_TITLES):
                slide_idx = 6 + i  # slides 7-11 are indices 6-10
                actual_title = get_slide_title(prs.slides[slide_idx])
                if actual_title == expected_title:
                    matched += 1
                else:
                    print(f"  INFO: Slide {slide_idx + 1} title mismatch: expected '{expected_title}', found '{actual_title}'")

            if matched == 5:
                print(f"PASS: Component 4 — All 5 shifted slides (7-11) have correct titles (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 4 — Only {matched}/5 shifted slides match")
        else:
            print(f"FAIL: Component 4 — Not enough slides ({num_slides}) to check shifted positions")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
