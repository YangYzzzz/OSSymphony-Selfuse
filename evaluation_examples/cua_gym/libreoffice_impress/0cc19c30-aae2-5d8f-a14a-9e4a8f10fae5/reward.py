"""
Reward Script: Duplicate slide 1 and place the copy immediately after slide 1.
Task ID: impstruct_003
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): Total slide count is 4
  Component 2 (0.3): Slide 2 title text matches Slide 1 title text ("Workshop Title")
  Component 3 (0.2): Slide 2 body content matches Slide 1 body content (duplicate)
  Component 4 (0.2): Original slides preserved at positions 3 and 4 ("Agenda", "Resources")
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impstruct_003'


def get_slide_texts(slide):
    """Extract all non-empty text strings from a slide's shapes."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return texts


def get_slide_title(slide):
    """Extract the title text from a slide, if present."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.shape_type is not None:
            # Check if shape is the title placeholder
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                if shape.placeholder_format.idx == 0:  # title placeholder
                    return shape.text_frame.text.strip()
    # Fallback: first non-empty text
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                return t
    return ""


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Total slide count is 4 (0.3 points)
    # Initial has 3 slides; golden should have 4 after duplication
    try:
        if num_slides == 4:
            print(f"PASS: Component 1 — Slide count is 4 (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 — Expected 4 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Early exit if slide count is wrong (can't verify positions)
    if num_slides < 4:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    slides = list(prs.slides)

    # Component 2: Slide 2 title matches Slide 1 title "Workshop Title" (0.3 points)
    # This checks the core duplication — the copy should have the same title
    try:
        title_1 = get_slide_title(slides[0])
        title_2 = get_slide_title(slides[1])
        if title_1 and title_2 and title_1 == title_2 and "Workshop" in title_1:
            print(f"PASS: Component 2 — Slide 2 title matches Slide 1: '{title_2}' (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 — Slide 1 title='{title_1}', Slide 2 title='{title_2}'")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 2 body content matches Slide 1 body content (0.2 points)
    # The duplicated slide should have all the same text content
    try:
        texts_1 = get_slide_texts(slides[0])
        texts_2 = get_slide_texts(slides[1])
        if texts_1 and texts_2 and texts_1 == texts_2:
            print(f"PASS: Component 3 — Slide 2 full text matches Slide 1 ({len(texts_1)} text items) (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 3 — Slide 1 texts={texts_1}, Slide 2 texts={texts_2}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Original slides preserved at positions 3 and 4 (0.2 points)
    # After duplication, "Agenda" should be slide 3 and "Resources" should be slide 4
    try:
        title_3 = get_slide_title(slides[2])
        title_4 = get_slide_title(slides[3])
        agenda_ok = "Agenda" in title_3 if title_3 else False
        resources_ok = "Resources" in title_4 if title_4 else False
        if agenda_ok and resources_ok:
            print(f"PASS: Component 4 — Slide 3='{title_3}', Slide 4='{title_4}' (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 4 — Slide 3='{title_3}' (expect 'Agenda'), Slide 4='{title_4}' (expect 'Resources')")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Main execution
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
