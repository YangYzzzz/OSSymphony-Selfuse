"""
Initial Setup: Duplicate slide 1 and place the copy immediately after slide 1.
Task ID: impstruct_003
Domain: libreoffice_impress

Creates a 3-slide presentation 'impstruct_003.pptx':
  Slide 1: Workshop Title
  Slide 2: Agenda
  Slide 3: Resources
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impstruct_003'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ----------------------------------------------------------------
    # Slide 1: Workshop Title (Title Slide layout)
    # ----------------------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Workshop Title"

    # Title formatting
    for run in slide1.shapes.title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    # Subtitle
    subtitle = slide1.placeholders[1]
    subtitle.text = "Advanced Data Visualization Techniques"
    for run in subtitle.text_frame.paragraphs[0].runs:
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0x4A, 0x4A, 0x4A)

    # Add a text box with date and presenter info
    txBox = slide1.shapes.add_textbox(Inches(3.5), Inches(5.0), Inches(6), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Presented by Dr. Elena Vasquez"
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    p2 = tf.add_paragraph()
    p2.text = "March 22, 2025 | Conference Room B"
    p2.alignment = PP_ALIGN.CENTER
    for run in p2.runs:
        run.font.size = Pt(14)
        run.font.italic = True
        run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # Background color for slide 1
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)

    # ----------------------------------------------------------------
    # Slide 2: Agenda (Title + Content layout)
    # ----------------------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    for run in slide2.shapes.title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    body = slide2.placeholders[1]
    tf_body = body.text_frame
    tf_body.clear()

    agenda_items = [
        "9:00 AM - Welcome and Introduction",
        "9:30 AM - Fundamentals of Chart Design",
        "10:15 AM - Interactive Dashboards with Python",
        "11:00 AM - Coffee Break",
        "11:15 AM - Real-Time Data Streaming Visualizations",
        "12:00 PM - Case Study: Financial Analytics Platform",
        "12:45 PM - Q&A and Wrap-Up",
    ]
    for i, item in enumerate(agenda_items):
        if i == 0:
            p = tf_body.paragraphs[0]
        else:
            p = tf_body.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ----------------------------------------------------------------
    # Slide 3: Resources (Title + Content layout)
    # ----------------------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Resources"
    for run in slide3.shapes.title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    body3 = slide3.placeholders[1]
    tf_body3 = body3.text_frame
    tf_body3.clear()

    resources = [
        "Matplotlib Documentation: matplotlib.org/stable/",
        "Plotly Dash Framework: dash.plotly.com",
        "D3.js Gallery: observablehq.com/@d3/gallery",
        "Seaborn Statistical Plots: seaborn.pydata.org",
        "Tableau Public: public.tableau.com",
        "Workshop GitHub Repository: github.com/evasquez/dataviz-workshop",
    ]
    for i, item in enumerate(resources):
        if i == 0:
            p = tf_body3.paragraphs[0]
        else:
            p = tf_body3.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0x2A, 0x5D, 0x9F)

    # Add contact info text box
    txBox3 = slide3.shapes.add_textbox(Inches(1), Inches(6.0), Inches(8), Inches(1))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p_contact = tf3.paragraphs[0]
    p_contact.text = "Contact: elena.vasquez@university.edu | Twitter: @dr_evasquez"
    p_contact.alignment = PP_ALIGN.LEFT
    for run in p_contact.runs:
        run.font.size = Pt(14)
        run.font.italic = True
        run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
