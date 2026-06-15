"""
Reward Script: Create a title slide for Introduction to Biology course
Task ID: impress_teach_001
Domain: libreoffice_impress
Scoring:
  Component 1 — Title text correct (0.25)
  Component 2 — Title font size 44pt (0.15)
  Component 3 — Title font bold (0.15)
  Component 4 — Subtitle text correct (0.25)
  Component 5 — Subtitle font size 20pt (0.20)
"""

import os

from pptx import Presentation
from pptx.util import Pt

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_001'

EXPECTED_TITLE = 'Introduction to Biology'
EXPECTED_SUBTITLE = 'Professor Sarah Chen | Fall 2025'
EXPECTED_TITLE_SIZE = Pt(44)   # 558800 EMU
EXPECTED_SUBTITLE_SIZE = Pt(20)  # 254000 EMU


def persist_app_state():
    """Save any unsaved LibreOffice edits before verifying."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def get_placeholder_text_and_runs(slide, placeholder_name_prefix):
    """Get full text and list of runs from a named placeholder."""
    for shape in slide.shapes:
        if shape.name.lower().startswith(placeholder_name_prefix.lower()):
            if shape.has_text_frame:
                full_text = shape.text_frame.text.strip()
                runs = []
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if (run.text or "").strip():
                            runs.append(run)
                return full_text, runs
    return None, []


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    if len(prs.slides) == 0:
        print("FAIL: No slides found in presentation")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]

    # --- Title checks ---
    title_text, title_runs = get_placeholder_text_and_runs(slide, "title")

    # Component 1: Title text is "Introduction to Biology" (0.25 points)
    try:
        if title_text and title_text == EXPECTED_TITLE:
            print(f"PASS: Component 1 - Title text is '{title_text}' (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 - Expected title '{EXPECTED_TITLE}', found '{title_text}'")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Title font size is 44pt (0.15 points)
    try:
        if title_runs:
            actual_size = title_runs[0].font.size
            if actual_size is not None and actual_size == EXPECTED_TITLE_SIZE:
                print(f"PASS: Component 2 - Title font size is 44pt ({actual_size} EMU) (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 2 - Expected title size {EXPECTED_TITLE_SIZE} EMU (44pt), found {actual_size}")
        else:
            print("FAIL: Component 2 - No title runs found to check font size")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Title font is bold (0.15 points)
    try:
        if title_runs:
            is_bold = title_runs[0].font.bold
            if is_bold is True:
                print(f"PASS: Component 3 - Title font is bold (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 3 - Expected title bold=True, found bold={is_bold}")
        else:
            print("FAIL: Component 3 - No title runs found to check bold")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # --- Subtitle checks ---
    subtitle_text, subtitle_runs = get_placeholder_text_and_runs(slide, "subtitle")

    # Component 4: Subtitle text is "Professor Sarah Chen | Fall 2025" (0.25 points)
    try:
        if subtitle_text and subtitle_text == EXPECTED_SUBTITLE:
            print(f"PASS: Component 4 - Subtitle text is '{subtitle_text}' (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 4 - Expected subtitle '{EXPECTED_SUBTITLE}', found '{subtitle_text}'")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Subtitle font size is 20pt (0.20 points)
    try:
        if subtitle_runs:
            actual_size = subtitle_runs[0].font.size
            if actual_size is not None and actual_size == EXPECTED_SUBTITLE_SIZE:
                print(f"PASS: Component 5 - Subtitle font size is 20pt ({actual_size} EMU) (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 5 - Expected subtitle size {EXPECTED_SUBTITLE_SIZE} EMU (20pt), found {actual_size}")
        else:
            print("FAIL: Component 5 - No subtitle runs found to check font size")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
