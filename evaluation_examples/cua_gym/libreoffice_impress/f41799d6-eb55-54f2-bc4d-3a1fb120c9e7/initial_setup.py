"""
Initial Setup: Build a 10-slide Executive Summary presentation with slide 5 as
an empty 'Performance Dashboard' ready for KPI content.
Task ID: impress_rp_050
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_050'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
    return slide


def add_blank_with_title(prs, title_text):
    """Add a slide with title only layout (layout index 5 = blank)."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    add_title_slide(prs, "Executive Summary", "FY2025 Annual Business Review\nPrepared by Strategic Planning Division")

    # Slide 2: Agenda
    add_content_slide(prs, "Agenda", [
        "1. Company Overview & Mission Update",
        "2. Financial Performance Highlights",
        "3. Market Analysis & Competitive Landscape",
        "4. Performance Dashboard",
        "5. Operational Metrics Deep Dive",
        "6. Strategic Initiatives for FY2026",
        "7. Risk Assessment & Mitigation",
        "8. Team Growth & Talent Pipeline",
        "9. Next Steps & Action Items",
    ])

    # Slide 3: Company Overview
    add_content_slide(prs, "Company Overview", [
        "Founded in 2012, NovaTech Solutions has grown to 2,400+ employees across 12 offices globally",
        "Annual revenue reached $487M in FY2025, representing 18% year-over-year growth",
        "Market share expanded to 14.3% in the enterprise software segment",
        "Customer retention rate stands at 94.7%, up from 91.2% last year",
        "Launched 3 new product lines: DataSync Pro, CloudBridge, and SecureVault",
    ])

    # Slide 4: Financial Highlights
    add_content_slide(prs, "Financial Performance Highlights", [
        "Revenue: $487.2M (target was $475M) - exceeded by 2.6%",
        "Gross Margin: 68.4% vs 65.1% prior year",
        "EBITDA: $112.3M representing 23.1% margin",
        "Operating Cash Flow: $98.7M, up 31% from FY2024",
        "R&D Investment: $73.1M (15% of revenue) focused on AI capabilities",
        "Customer Acquisition Cost reduced by 12% to $2,340 per enterprise client",
    ])

    # Slide 5: Performance Dashboard - EMPTY (this is the task target)
    slide5 = add_blank_with_title(prs, "Performance Dashboard")
    # Leave slide 5 intentionally empty - the agent task is to build the KPI dashboard here

    # Slide 6: Operational Metrics
    add_content_slide(prs, "Operational Metrics Deep Dive", [
        "Support ticket resolution time: 4.2 hours (down from 6.8 hours)",
        "System uptime: 99.97% across all production environments",
        "Deployment frequency: 47 releases per quarter (CI/CD pipeline)",
        "Mean time to recovery: 12 minutes for P1 incidents",
        "API response time p95: 145ms (target < 200ms)",
    ])

    # Slide 7: Strategic Initiatives
    add_content_slide(prs, "Strategic Initiatives FY2026", [
        "AI-Powered Analytics Engine - Expected launch Q2 2026",
        "APAC Market Expansion - New offices in Singapore and Tokyo",
        "Enterprise Security Suite - SOC2 Type II certification in progress",
        "Partner Ecosystem Growth - Target 150+ certified integration partners",
        "Sustainability Program - Carbon neutral operations by Q4 2026",
    ])

    # Slide 8: Risk Assessment
    add_content_slide(prs, "Risk Assessment & Mitigation", [
        "Competitive pressure from 3 new market entrants in enterprise segment",
        "Supply chain disruptions affecting hardware procurement timelines",
        "Regulatory changes in EU data sovereignty requirements (GDPR extensions)",
        "Talent retention in engineering - 8.2% attrition rate vs 6.5% target",
        "Currency exposure in international revenue streams (EUR, GBP, JPY)",
    ])

    # Slide 9: Team & Talent
    add_content_slide(prs, "Team Growth & Talent Pipeline", [
        "Headcount grew from 1,980 to 2,412 employees (+21.8%)",
        "Engineering team: 890 members across 42 squads",
        "Average tenure increased to 3.4 years from 2.9 years",
        "Internal promotion rate: 28% of all role fills",
        "Diversity hiring: 47% of new hires from underrepresented groups",
    ])

    # Slide 10: Next Steps
    add_content_slide(prs, "Next Steps & Action Items", [
        "Finalize FY2026 budget allocation by March 31, 2026",
        "Complete Q1 OKR alignment sessions across all business units",
        "Launch customer advisory board with top 20 enterprise clients",
        "Submit SOC2 Type II audit documentation by April 15, 2026",
        "Quarterly business review scheduled for June 12, 2026",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
