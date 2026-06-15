"""
Reward Script: KPI Dashboard on Slide 5
Task ID: impress_rp_050
Domain: libreoffice_impress
Scoring:
  Component 1: Gauge arc shapes (green + gray) in top-left — 0.25 pts
  Component 2: "78%" text label in gauge area — 0.15 pts
  Component 3: 4 bar rectangles (Q1-Q4) in top-right — 0.25 pts
  Component 4: Arrow shape colored #27AE60 in bottom-left — 0.15 pts
  Component 5: KPI text box with 3 lines in bottom-right — 0.20 pts
"""

import os
import xml.etree.ElementTree as ET
import zipfile

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_050'

def persist_app_state(domain):
    """Save any unsaved GUI edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            import time
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify KPI dashboard on slide 5 with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # 0-indexed → slide 5
    shapes = list(slide.shapes)

    # Helper: get line color from shape XML
    def get_line_color_from_xml(shape):
        """Extract line/outline color from shape XML."""
        try:
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            xml_str = shape._element.xml
            root = ET.fromstring(xml_str)
            ln = root.find('.//' + '{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
            if ln is not None:
                srgb = ln.find('.//' + '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                if srgb is not None:
                    return srgb.get('val', '').upper()
            return None
        except Exception:
            return None

    def get_fill_color(shape):
        """Get solid fill color from shape."""
        try:
            fill = shape.fill
            if fill.type is not None and fill.type == 1:  # SOLID
                return str(fill.fore_color.rgb).upper()
        except Exception:
            pass
        return None

    def is_arc_shape(shape):
        """Check if shape is an arc (preset geometry 'arc')."""
        try:
            xml_str = shape._element.xml
            return 'prst="arc"' in xml_str
        except Exception:
            return False

    # Slide midpoints for quadrant classification
    slide_mid_x = prs.slide_width // 2
    slide_mid_y = (prs.slide_height + 914400) // 2  # offset for title area

    # ---------------------------------------------------------------
    # Component 1: Gauge arc shapes (green + gray) in top-left (0.25)
    # Task requires: arc shapes forming semicircle gauge,
    #   green (#27AE60) and gray (#E0E0E0) arcs
    # ---------------------------------------------------------------
    try:
        arc_shapes = [s for s in shapes if is_arc_shape(s)]
        green_arc_found = False
        gray_arc_found = False

        for arc in arc_shapes:
            line_color = get_line_color_from_xml(arc)
            fill_color = get_fill_color(arc)
            color = line_color or fill_color
            if color:
                if color == '27AE60':
                    green_arc_found = True
                elif color == 'E0E0E0':
                    gray_arc_found = True

        if len(arc_shapes) >= 2 and green_arc_found and gray_arc_found:
            print(f"PASS: Component 1 — Found {len(arc_shapes)} arc shapes with green (#27AE60) and gray (#E0E0E0) gauge (0.25 pts)")
            total_score += 0.25
        elif len(arc_shapes) >= 2 and (green_arc_found or gray_arc_found):
            print(f"PARTIAL: Component 1 — Found arcs but only one color matched (0.10 pts)")
            total_score += 0.10
        elif len(arc_shapes) >= 1:
            print(f"PARTIAL: Component 1 — Found {len(arc_shapes)} arc(s) but colors don't match. green={green_arc_found}, gray={gray_arc_found}")
            total_score += 0.05
        else:
            print(f"FAIL: Component 1 — No arc shapes found on slide 5. Found {len(shapes)} total shapes.")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # ---------------------------------------------------------------
    # Component 2: "78%" text label in gauge area (0.15)
    # ---------------------------------------------------------------
    try:
        pct_text_found = False
        for shape in shapes:
            if hasattr(shape, 'text') and shape.text:
                text_stripped = shape.text.strip()
                if '78%' in text_stripped or '78 %' in text_stripped:
                    pct_text_found = True
                    break

        if pct_text_found:
            print(f"PASS: Component 2 — Found '78%' text label on slide 5 (0.15 pts)")
            total_score += 0.15
        else:
            all_texts = [s.text.strip() for s in shapes if hasattr(s, 'text') and s.text]
            print(f"FAIL: Component 2 — '78%' text not found. Texts on slide: {all_texts}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # ---------------------------------------------------------------
    # Component 3: 4 bar rectangles (Q1-Q4) in top-right (0.25)
    # Task requires: 4 rectangular bars with Q1-Q4 labels
    # ---------------------------------------------------------------
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        rect_shapes = []
        for s in shapes:
            try:
                # Only count non-placeholder, non-textbox rectangles with solid fill
                # (bars should have solid fill colors, not be text containers)
                if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    xml_str = s._element.xml
                    if 'prst="rect"' in xml_str:
                        fill_color = get_fill_color(s)
                        if fill_color is not None:
                            # Has solid fill — likely a bar, not a text container
                            rect_shapes.append(s)
            except Exception:
                pass

        # Check for Q1-Q4 labels
        quarter_labels = set()
        for shape in shapes:
            if hasattr(shape, 'text') and shape.text:
                t = shape.text.strip()
                if t in ('Q1', 'Q2', 'Q3', 'Q4'):
                    quarter_labels.add(t)

        num_bars = len(rect_shapes)
        num_labels = len(quarter_labels)

        if num_bars >= 4 and num_labels >= 4:
            print(f"PASS: Component 3 — Found {num_bars} bar rectangles and labels {quarter_labels} (0.25 pts)")
            total_score += 0.25
        elif num_bars >= 4 and num_labels >= 2:
            print(f"PARTIAL: Component 3 — Found {num_bars} bars but only {num_labels} labels: {quarter_labels} (0.15 pts)")
            total_score += 0.15
        elif num_bars >= 2:
            print(f"PARTIAL: Component 3 — Found {num_bars} bars and {num_labels} labels (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 3 — Found {num_bars} rectangles and {num_labels} quarter labels")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # ---------------------------------------------------------------
    # Component 4: Arrow shape colored #27AE60 (0.15)
    # Task requires: trend arrow pointing up-right, green (#27AE60)
    # ---------------------------------------------------------------
    try:
        arrow_found = False
        arrow_color_correct = False

        for shape in shapes:
            xml_str = shape._element.xml
            # Check for arrow preset geometries
            is_arrow = any(kw in xml_str for kw in [
                'prst="upArrow"', 'prst="rightArrow"',
                'prst="upDownArrow"', 'prst="bentUpArrow"',
                'prst="notchedRightArrow"', 'prst="stripedRightArrow"',
                'name="Up Arrow'
            ])
            if is_arrow:
                arrow_found = True
                color = get_fill_color(shape)
                if color and color == '27AE60':
                    arrow_color_correct = True
                    break

        if arrow_found and arrow_color_correct:
            print(f"PASS: Component 4 — Arrow shape with fill #27AE60 found (0.15 pts)")
            total_score += 0.15
        elif arrow_found:
            print(f"PARTIAL: Component 4 — Arrow shape found but color mismatch (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 4 — No arrow shape found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # ---------------------------------------------------------------
    # Component 5: KPI text box with 3 lines in bottom-right (0.20)
    # Task requires: summary text box with 3 lines of KPI metrics
    # ---------------------------------------------------------------
    try:
        kpi_textbox_found = False

        for shape in shapes:
            if not (hasattr(shape, 'text_frame') and hasattr(shape, 'text')):
                continue
            text = shape.text.strip()
            # KPI text box should have multiple lines with metric-like content
            lines = [l.strip() for l in text.split('\n') if l.strip()]
            if len(lines) >= 3:
                # Check if lines contain KPI-like content (numbers, percentages, metrics)
                has_numbers = sum(1 for l in lines if any(c.isdigit() for c in l))
                if has_numbers >= 2:
                    # Verify this is in the bottom portion of the slide (not the title)
                    if shape.top > 914400:  # below title area
                        kpi_textbox_found = True
                        print(f"  KPI lines found: {lines}")
                        break

        if kpi_textbox_found:
            print(f"PASS: Component 5 — KPI text box with 3+ metric lines found (0.20 pts)")
            total_score += 0.20
        else:
            # Check for partial: any multi-line text box with numbers
            partial_found = False
            for shape in shapes:
                if hasattr(shape, 'text') and shape.text:
                    lines = [l.strip() for l in shape.text.split('\n') if l.strip()]
                    if len(lines) >= 2 and any(any(c.isdigit() for c in l) for l in lines):
                        if shape.top > 914400:
                            partial_found = True
                            break
            if partial_found:
                print(f"PARTIAL: Component 5 — Found multi-line text but fewer than 3 KPI lines (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 5 — No KPI text box with 3 metric lines found")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
