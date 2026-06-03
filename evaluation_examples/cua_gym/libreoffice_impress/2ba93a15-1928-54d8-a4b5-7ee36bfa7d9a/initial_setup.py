"""
Initial Setup: Competitive Analysis presentation with 8 slides, slide 6 has no chart.
Task ID: impress_exec_091
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_091'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
    return slide


def add_title_only_slide(prs, title_text):
    """Add a slide with only a title (layout 5=blank with manual title, or layout 6=title only)."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)
    return slide


def create_initial():
    prs = Presentation()

    # Slide 1: Title Slide
    add_title_slide(prs, "Competitive Analysis Report", "Q1 2025 Strategic Review\nPrepared by Market Intelligence Team")

    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "Overall market share increased 3.2% year-over-year",
        "Key competitor launched aggressive pricing strategy in Q4 2024",
        "Our product innovation pipeline remains strongest in sector",
        "Customer satisfaction scores lead industry at 92%",
        "Brand perception improved significantly in under-35 demographic",
    ])

    # Slide 3: Market Overview
    add_content_slide(prs, "Market Overview", [
        "Total addressable market: $4.7B (up 8% from 2024)",
        "Our market share: 23.4% (+1.8pp YoY)",
        "Top competitor market share: 21.1% (+0.3pp YoY)",
        "Three new entrants in niche segments",
        "Regulatory changes expected H2 2025 may affect pricing",
    ])

    # Slide 4: Our Product Portfolio
    add_content_slide(prs, "Our Product Portfolio", [
        "Flagship product: ProSuite Enterprise (revenue +18% YoY)",
        "Mid-tier: ProSuite Standard (steady growth at 7%)",
        "New launch: ProSuite AI Assistant (beta since Jan 2025)",
        "Legacy support maintained for ProSuite Classic users",
        "R&D investment: $142M across 6 product lines",
    ])

    # Slide 5: Competitor Landscape
    add_content_slide(prs, "Competitor Landscape", [
        "Acme Corp: Strongest in pricing, recent brand refresh",
        "TechVantage: Innovation-driven but weak customer support",
        "GlobalReach: Strong brand loyalty in European markets",
        "DataFirst: Emerging player with AI-first approach",
        "Market consolidation expected with 2-3 acquisitions predicted",
    ])

    # Slide 6: Head-to-Head Comparison (NO CHART - this is what the agent must add)
    slide6 = add_title_only_slide(prs, "Head-to-Head Comparison")
    # Add a subtitle/description text box
    txBox = slide6.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "This slide needs a radar chart showing competitive positioning across key dimensions."
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(14)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

    # Slide 7: Strategic Recommendations
    add_content_slide(prs, "Strategic Recommendations", [
        "Invest in support infrastructure to maintain competitive advantage",
        "Develop pricing tiers to counter Acme Corp's aggressive strategy",
        "Accelerate AI product launch to capture innovation narrative",
        "Strengthen brand marketing in key growth demographics",
        "Explore strategic partnerships in European markets",
    ])

    # Slide 8: Next Steps
    add_content_slide(prs, "Next Steps & Timeline", [
        "Q2 2025: Launch competitive pricing pilot program",
        "Q2 2025: Complete AI Assistant general availability",
        "Q3 2025: Execute brand awareness campaign",
        "Q3 2025: Expand enterprise support team by 40%",
        "Q4 2025: Full competitive review and strategy refresh",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
