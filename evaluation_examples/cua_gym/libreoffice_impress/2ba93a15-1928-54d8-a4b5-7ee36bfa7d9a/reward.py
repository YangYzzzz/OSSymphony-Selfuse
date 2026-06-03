"""
Reward Script: Add radar chart on slide 6 with competitive analysis data
Task ID: impress_exec_091
Domain: libreoffice_impress
Scoring:
  - Component 1: Chart exists on slide 6 (0.20)
  - Component 2: Chart type is radar (0.15)
  - Component 3: Chart title is 'Competitive Positioning' (0.15)
  - Component 4: 5 correct categories (0.20)
  - Component 5: 'Us' series data correct (0.15)
  - Component 6: 'Competitor' series data correct (0.15)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_091'


def persist_app_state(domain: str):
    """Save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_series_name(series):
    """Extract series name from XML since python-pptx doesn't expose it directly."""
    try:
        ns = '{http://schemas.openxmlformats.org/drawingml/2006/chart}'
        tx = series._element.find(f'.//{ns}tx')
        if tx is not None:
            v = tx.find(f'.//{ns}v')
            if v is not None:
                return v.text
    except Exception:
        pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Missing python-pptx library: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"PRECONDITION FAIL: Expected at least 6 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[5]  # Slide 6 (0-indexed)

    # Find chart shape on slide 6
    chart_shape = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            chart_shape = shape
            break

    # Component 1: Chart exists on slide 6 (0.20 points)
    try:
        if chart_shape is not None:
            print(f"PASS: Component 1 -- Chart found on slide 6 (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 -- No chart found on slide 6")
            # No chart means all subsequent checks will fail too
            final_score = min(total_score, 1.0)
            print(f"\nScore: {total_score}/1.0")
            print(f"REWARD: {final_score}")
            return final_score
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    chart = chart_shape.chart

    # Component 2: Chart type is radar (0.15 points)
    # Radar types: RADAR (81), RADAR_FILLED (82), RADAR_MARKERS (83)
    try:
        chart_type_val = chart.chart_type
        # Check if it's any radar type (enum values 81, 82, 83)
        radar_types = {81, 82, 83}
        if int(chart_type_val) in radar_types:
            print(f"PASS: Component 2 -- Chart is radar type (value={int(chart_type_val)}) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 -- Chart type is {chart_type_val} (value={int(chart_type_val)}), expected radar (81/82/83)")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Chart title is 'Competitive Positioning' (0.15 points)
    try:
        if chart.has_title:
            title_text = chart.chart_title.text_frame.text.strip()
            if title_text.lower() == 'competitive positioning':
                print(f"PASS: Component 3 -- Chart title is '{title_text}' (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 3 -- Chart title is '{title_text}', expected 'Competitive Positioning'")
        else:
            print(f"FAIL: Component 3 -- Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: 5 correct categories (0.20 points)
    expected_categories = ['Product', 'Price', 'Support', 'Innovation', 'Brand']
    try:
        plot = chart.plots[0]
        actual_cats = list(plot.categories)
        if len(actual_cats) == 5:
            # Case-insensitive comparison
            cats_match = all(
                a.strip().lower() == e.lower()
                for a, e in zip(actual_cats, expected_categories)
            )
            if cats_match:
                print(f"PASS: Component 4 -- Categories match: {actual_cats} (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 4 -- Categories mismatch: {actual_cats} vs expected {expected_categories}")
        else:
            print(f"FAIL: Component 4 -- Expected 5 categories, found {len(actual_cats)}: {actual_cats}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: 'Us' series data correct (0.15 points)
    expected_us = [9.0, 6.0, 8.0, 9.0, 7.0]
    try:
        plot = chart.plots[0]
        series_list = list(plot.series)
        if len(series_list) < 2:
            print(f"FAIL: Component 5 -- Expected 2 series, found {len(series_list)}")
        else:
            # Find the 'Us' series (check both series)
            us_found = False
            for series in series_list:
                sname = get_series_name(series)
                vals = list(series.values)
                if sname and sname.strip().lower() == 'us':
                    if vals == expected_us:
                        print(f"PASS: Component 5 -- 'Us' series values correct: {vals} (0.15 pts)")
                        total_score += 0.15
                    else:
                        print(f"FAIL: Component 5 -- 'Us' series values: {vals}, expected {expected_us}")
                    us_found = True
                    break
            if not us_found:
                # Try matching by values if names don't match
                for series in series_list:
                    vals = list(series.values)
                    if vals == expected_us:
                        sname = get_series_name(series) or 'unnamed'
                        print(f"PASS: Component 5 -- Series '{sname}' has correct 'Us' values: {vals} (0.15 pts)")
                        total_score += 0.15
                        us_found = True
                        break
                if not us_found:
                    names = [get_series_name(s) for s in series_list]
                    print(f"FAIL: Component 5 -- No series with 'Us' data found. Series names: {names}")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # Component 6: 'Competitor' series data correct (0.15 points)
    expected_comp = [7.0, 8.0, 5.0, 6.0, 8.0]
    try:
        plot = chart.plots[0]
        series_list = list(plot.series)
        if len(series_list) < 2:
            print(f"FAIL: Component 6 -- Expected 2 series, found {len(series_list)}")
        else:
            comp_found = False
            for series in series_list:
                sname = get_series_name(series)
                vals = list(series.values)
                if sname and sname.strip().lower() == 'competitor':
                    if vals == expected_comp:
                        print(f"PASS: Component 6 -- 'Competitor' series values correct: {vals} (0.15 pts)")
                        total_score += 0.15
                    else:
                        print(f"FAIL: Component 6 -- 'Competitor' series values: {vals}, expected {expected_comp}")
                    comp_found = True
                    break
            if not comp_found:
                for series in series_list:
                    vals = list(series.values)
                    if vals == expected_comp:
                        sname = get_series_name(series) or 'unnamed'
                        print(f"PASS: Component 6 -- Series '{sname}' has correct 'Competitor' values: {vals} (0.15 pts)")
                        total_score += 0.15
                        comp_found = True
                        break
                if not comp_found:
                    names = [get_series_name(s) for s in series_list]
                    print(f"FAIL: Component 6 -- No series with 'Competitor' data found. Series names: {names}")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
