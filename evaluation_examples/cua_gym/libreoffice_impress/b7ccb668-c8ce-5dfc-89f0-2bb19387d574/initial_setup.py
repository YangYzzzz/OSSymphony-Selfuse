"""
Initial Setup: Remove transition from slide 5
Task ID: impress_tm_005
Domain: libreoffice_impress

Creates a 10-slide marketing plan presentation. Slide 5 ("Target Audience")
has a Checkerboard transition with 3.0s duration.
"""

import os
import shlex
import subprocess
import time
import zipfile
import shutil
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_005'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=None):
    """Helper to add a text box with formatting."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)
    return txBox


def add_bullet_points(slide, left, top, width, height, items, font_size=14):
    """Add bulleted text with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(font_size)
    return txBox


def create_initial():
    prs = Presentation()

    slide_data = [
        {
            'layout': 0,
            'title': 'Q3 2025 Marketing Plan',
            'subtitle': 'Prepared by the Growth Marketing Team\nAcme Technologies Inc.',
        },
        {
            'layout': 1,
            'title': 'Executive Summary',
            'bullets': [
                'Expand market presence in APAC region by 35%',
                'Launch integrated digital campaign across 6 channels',
                'Projected ROI of 4.2x on total marketing spend of $2.8M',
                'Key focus: enterprise SaaS vertical and mid-market segment',
                'Partnership with 12 regional distributors confirmed',
            ],
        },
        {
            'layout': 1,
            'title': 'Market Analysis',
            'bullets': [
                'Total addressable market: $14.7B (growing 18% YoY)',
                'Primary competitors: TechVista, CloudPeak, DataBridge',
                'Customer acquisition cost decreased 22% since Q1',
                'Net promoter score improved from 42 to 61',
                'Market share increased from 8.3% to 11.7%',
            ],
        },
        {
            'layout': 1,
            'title': 'Campaign Strategy',
            'bullets': [
                'Phase 1 (Jul): Brand awareness via LinkedIn and Google Ads',
                'Phase 2 (Aug): Content marketing with 15 whitepapers',
                'Phase 3 (Sep): Lead generation with webinar series',
                'Budget allocation: Digital 45%, Events 30%, PR 25%',
                'A/B testing on all landing pages and email sequences',
            ],
        },
        {
            'layout': 1,
            'title': 'Target Audience',
            'bullets': [
                'Primary: CTOs and VP Engineering at companies with 500-5000 employees',
                'Secondary: IT Directors in financial services and healthcare',
                'Tertiary: DevOps leads evaluating infrastructure modernization',
                'Geographic focus: Singapore, Tokyo, Sydney, Mumbai',
                'Decision-making cycle: 3-6 months average',
                'Average deal size: $85,000 annually',
            ],
        },
        {
            'layout': 1,
            'title': 'Budget Breakdown',
            'bullets': [
                'Digital advertising: $1,260,000 (45%)',
                'Trade shows and events: $840,000 (30%)',
                'Public relations and communications: $700,000 (25%)',
                'Contingency reserve: $140,000 (5% of total)',
                'Agency retainer fees: $320,000 included in digital',
            ],
        },
        {
            'layout': 1,
            'title': 'Timeline & Milestones',
            'bullets': [
                'July 1: Campaign launch across all digital channels',
                'July 15: First webinar - "Cloud Migration Best Practices"',
                'August 10: Singapore Tech Summit sponsorship',
                'August 28: Mid-campaign performance review',
                'September 15: Tokyo Enterprise Forum keynote',
                'September 30: Q3 campaign wrap-up and analysis',
            ],
        },
        {
            'layout': 1,
            'title': 'Key Performance Indicators',
            'bullets': [
                'Marketing qualified leads (MQLs): Target 2,400',
                'Sales qualified leads (SQLs): Target 480',
                'Website traffic increase: 65% over Q2 baseline',
                'Email open rate: maintain above 28%',
                'Cost per lead: reduce to below $45',
                'Pipeline contribution: $12M in new opportunities',
            ],
        },
        {
            'layout': 1,
            'title': 'Team & Resources',
            'bullets': [
                'Project lead: Sarah Chen, VP Marketing',
                'Digital team: Marcus Johnson, Priya Patel, Takeshi Yamamoto',
                'Content: Elena Rodriguez, James O\'Brien',
                'Analytics: David Kim, Aisha Okafor',
                'External agency: BrightWave Digital (creative + media buying)',
            ],
        },
        {
            'layout': 1,
            'title': 'Next Steps',
            'bullets': [
                'Finalize creative assets by June 25',
                'Complete vendor agreements for all event sponsorships',
                'Set up tracking dashboards in HubSpot and GA4',
                'Schedule bi-weekly stakeholder review meetings',
                'Distribute campaign playbook to regional teams',
            ],
        },
    ]

    for i, sd in enumerate(slide_data):
        layout = prs.slide_layouts[sd['layout']]
        slide = prs.slides.add_slide(layout)

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = sd['title']

        if sd['layout'] == 0:
            # Title slide with subtitle
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = sd.get('subtitle', '')
        elif 'bullets' in sd:
            # Content slide with bullet points
            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for j, bullet in enumerate(sd['bullets']):
                    if j == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
                    for run in p.runs:
                        run.font.size = Pt(16)

    # Save the presentation first
    prs.save(OUTPUT)
    print(f'Initial presentation created: {OUTPUT}')

    # Now add Checkerboard transition to slide 5 via XML manipulation
    # python-pptx doesn't support transitions, so we edit the XML directly
    tmp_path = f'{WORKDIR}/{TASK_ID}_tmp.pptx'
    shutil.copy(OUTPUT, tmp_path)

    # Read slide5.xml, add transition element
    with zipfile.ZipFile(tmp_path, 'r') as zin:
        with zipfile.ZipFile(OUTPUT, 'w') as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == 'ppt/slides/slide5.xml':
                    # Parse and add transition
                    root = ET.fromstring(data)
                    # Define namespaces
                    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                    ns_p14 = 'http://schemas.microsoft.com/office/powerpoint/2010/main'
                    ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

                    # Register namespaces to preserve them
                    namespaces = {
                        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                        'r': ns_r,
                        'p': ns_p,
                    }
                    for prefix, uri in namespaces.items():
                        ET.register_namespace(prefix, uri)

                    # Remove existing transition if any
                    for tr in root.findall(f'{{{ns_p}}}transition'):
                        root.remove(tr)

                    # Create transition element: Checkerboard with 3.0s duration
                    # Duration in ms: 3000
                    transition = ET.SubElement(root, f'{{{ns_p}}}transition')
                    transition.set('spd', 'slow')  # slow = ~3s
                    transition.set('advClick', '1')

                    # Add checker child element
                    checker = ET.SubElement(transition, f'{{{ns_p}}}checker')
                    checker.set('dir', 'horz')

                    data = ET.tostring(root, encoding='unicode', xml_declaration=True).encode('utf-8')

                zout.writestr(item, data)

    # Remove temp file
    os.remove(tmp_path)
    print('Added Checkerboard transition to slide 5')

    # Launch LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
