"""
FINAL REWARD SCRIPT - SUCCESS
Task: Could you help me put the 'Conclusion' slide all the way at the end of my presentation? It's currently not in the right spot.
Generated: 2025-08-07 12:09:49
Status: success
Model: o4-mini
Total Steps: 1
"""

import os
from pptx import Presentation

def verify_conclusion_position(file_path):
    """
    Verifies that the presentation at file_path has a 'Conclusion' slide and that it is the last slide.
    Scoring breakdown:
      - 0.2 points if the file exists
      - 0.2 points if the file loads successfully
      - 0.3 points if a 'Conclusion' slide is detected
      - 0.3 points if the 'Conclusion' slide is at the end
    Total: 1.0 points maximum
    """
    print("Starting verification of 'Conclusion' slide position...")
    total_score = 0.0

    # 1. Check file existence (0.2)
    if os.path.exists(file_path):
        print("✓ File exists (0.2)")
        total_score += 0.2
    else:
        print(f"✗ File not found: {file_path}")
        print(f"REWARD: {total_score}")
        return total_score

    # 2. Load presentation (0.2)
    try:
        prs = Presentation(file_path)
        num_slides = len(prs.slides)
        print(f"✓ Presentation loaded successfully with {num_slides} slides (0.2)")
        total_score += 0.2
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        print(f"REWARD: {total_score}")
        return total_score

    # 3. Detect 'Conclusion' slide(s) (0.3)
    conclusion_indices = []
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text and "conclusion" in text.lower():
                    conclusion_indices.append(idx)
                    print(f"✓ Found 'Conclusion' in slide {idx+1}")
                    break
    if conclusion_indices:
        print(f"✓ 'Conclusion' slide(s) detected at positions: {[i+1 for i in conclusion_indices]} (0.3)")
        total_score += 0.3
    else:
        print("✗ No 'Conclusion' slide detected (0.0 for detection step)")

    # 4. Verify 'Conclusion' is last slide (0.3)
    last_index = num_slides - 1
    if conclusion_indices and last_index in conclusion_indices:
        print(f"✓ 'Conclusion' slide is at the end (slide {last_index+1}) (0.3)")
        total_score += 0.3
    else:
        print(f"✗ 'Conclusion' slide is not at the end (expected slide {last_index+1}) (0.0 for position step)")

    # Cap score at 1.0
    final_score = min(total_score, 1.0)
    final_score = round(final_score, 3)
    print(f"Final score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score

if __name__ == '__main__':
    file_path = '/home/user/could_you_help_me_put_the_conclusion_slide_all_the_way_at_the_end_of_my_presentation_its_currently_n.pptx'
    verify_conclusion_position(file_path)

