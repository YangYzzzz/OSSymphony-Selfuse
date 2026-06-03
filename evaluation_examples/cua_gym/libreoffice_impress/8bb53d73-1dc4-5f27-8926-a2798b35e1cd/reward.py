"""
FINAL REWARD SCRIPT - SUCCESS
Task: Insert the file name field in the header and align it center.
Generated: 2025-10-17 13:59:59
Status: success
Model: azure-o3
Total Steps: 14
"""

import os
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

"""
Reward verification script for the task:
"Insert the file name field in the header and align it center."

Scoring logic (progressive):
  • 0.6 weight – the slide actually shows the file name in its header area
  • 0.4 weight – that header text is centre-aligned

A perfect score (1.0) is awarded only when EVERY slide contains the
file name text AND that text is centred on EVERY slide.  Partial fulfilment
receives proportional credit.  No points are given for natural conditions
(file exists, slide exists, etc.).
"""


def verify_filename_header(file_path: str) -> float:
    """Returns a reward score between 0.0 and 1.0 for the task."""

    # Derive several acceptable patterns of the file name that might appear
    basename = os.path.basename(file_path)                       # full name e.g. demo.pptx
    base_no_golden = basename.replace("_golden", "")            # strip training suffix if still present
    name_no_ext = os.path.splitext(base_no_golden)[0]           # strip .pptx
    patterns = {basename, base_no_golden, name_no_ext}          # unique, non-empty patterns
    patterns = {p for p in patterns if p}                       # ensure no blanks

    # Load presentation (no points for just loading – prerequisite)
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Could not load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    if len(prs.slides) == 0:
        print("✗ Presentation contains no slides")
        print("REWARD: 0.0")
        return 0.0

    slides_total = len(prs.slides)
    slides_with_name = 0              # slides where file name text is present
    slides_with_name_center = 0       # slides where that text is centre aligned

    for idx, slide in enumerate(prs.slides, start=1):
        found_name = False
        found_center = False

        # Iterate through shapes that have text frames
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue

            text = shape.text_frame.text or ""
            if any(p in text for p in patterns):
                found_name = True
                align = shape.text_frame.paragraphs[0].alignment
                if align == PP_ALIGN.CENTER:
                    found_center = True
                # keep looking, there could be other occurrences

        # Debug per-slide result
        print(f"Slide {idx}: name_found={found_name}, centered={found_center}")

        if found_name:
            slides_with_name += 1
            if found_center:
                slides_with_name_center += 1

    # Calculate progressive score
    ratio_name = slides_with_name / slides_total         # proportion of slides with name present
    ratio_center = slides_with_name_center / slides_total  # proportion with centred alignment

    score = 0.6 * ratio_name + 0.4 * ratio_center

    # Debug summary
    print("-" * 40)
    print(f"Slides total           : {slides_total}")
    print(f"Slides with file name  : {slides_with_name}")
    print(f"Slides centred         : {slides_with_name_center}")
    print(f"Name ratio   (60%)     : {ratio_name:.2f}")
    print(f"Center ratio (40%)     : {ratio_center:.2f}")
    print(f"Computed score         : {score}")

    final_score = min(1.0, score)  # safety cap
    print(f"REWARD: {final_score}")
    return final_score


if __name__ == "__main__":
    FILE_PATH = "/home/user/insert_the_file_name_field_in_the_header_and_align_it_center.pptx"
    verify_filename_header(FILE_PATH)

