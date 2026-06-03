"""
Reward Script: Discussion Questions on Slide 9
Task ID: impress_teach_072
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): All 4 numbered questions present on slide 9
  Component 2 (0.3): Font size is 22pt (279400 EMU) for question text
  Component 3 (0.3): Line spacing is 2.0 for question paragraphs
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_072'

EXPECTED_QUESTIONS = [
    '1. How does climate affect biodiversity?',
    '2. What role do humans play in ecosystem change?',
    '3. Can we reverse the damage to coral reefs?',
    '4. What policies would you recommend?',
]

TARGET_FONT_SIZE = 279400  # 22pt in EMU (22 * 12700)
TARGET_LINE_SPACING = 2.0


def persist_app_state(domain):
    """Save any unsaved LibreOffice changes before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
    except ImportError:
        print("CRITICAL: python-pptx not installed")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 9 slides
    if len(prs.slides) < 9:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 9")
        print("REWARD: 0.0")
        return 0.0

    slide9 = prs.slides[8]  # 0-indexed

    # Collect all non-empty text from shapes on slide 9 that are NOT the pre-existing shapes.
    # We need to find the discussion questions text. We look for paragraphs containing numbered questions.
    # The initial slide 9 has: an empty Title placeholder and a "Discussion Time" textbox.
    # The golden adds a new textbox with the 4 questions.
    # We search all text shapes for the question content.

    all_question_paragraphs = []  # list of (text, font_size, line_spacing)

    for shape in slide9.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            # Check if this paragraph matches any expected question
            if text and any(text == q for q in EXPECTED_QUESTIONS):
                # Get font size from first non-empty run
                font_size = None
                for run in para.runs:
                    if (run.text or "").strip():
                        font_size = run.font.size
                        break
                line_spacing = para.line_spacing
                all_question_paragraphs.append((text, font_size, line_spacing))

    # Component 1: All 4 numbered questions present on slide 9 (0.4 points)
    try:
        found_texts = [t for t, _, _ in all_question_paragraphs]
        matched_count = sum(1 for q in EXPECTED_QUESTIONS if q in found_texts)
        if matched_count == 4:
            print(f"PASS: Component 1 -- All 4 discussion questions found on slide 9 (0.4 pts)")
            total_score += 0.4
        elif matched_count > 0:
            partial = round(0.4 * matched_count / 4, 2)
            print(f"PARTIAL: Component 1 -- {matched_count}/4 questions found (partial {partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- No discussion questions found on slide 9. Found texts: {found_texts}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Font size is 22pt (279400 EMU) for all question runs (0.3 points)
    try:
        if not all_question_paragraphs:
            print("FAIL: Component 2 -- No question paragraphs to check font size")
        else:
            correct_size_count = 0
            for text, font_size, _ in all_question_paragraphs:
                if font_size is not None and abs(font_size - TARGET_FONT_SIZE) < 1000:
                    correct_size_count += 1
                else:
                    print(f"  Font size mismatch for '{text[:40]}...': expected {TARGET_FONT_SIZE}, got {font_size}")

            if correct_size_count == len(all_question_paragraphs) and correct_size_count > 0:
                print(f"PASS: Component 2 -- All {correct_size_count} questions have 22pt font (0.3 pts)")
                total_score += 0.3
            elif correct_size_count > 0:
                partial = round(0.3 * correct_size_count / len(all_question_paragraphs), 2)
                print(f"PARTIAL: Component 2 -- {correct_size_count}/{len(all_question_paragraphs)} correct font size (partial {partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 2 -- No questions have correct 22pt font size")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Line spacing is 2.0 for all question paragraphs (0.3 points)
    try:
        if not all_question_paragraphs:
            print("FAIL: Component 3 -- No question paragraphs to check line spacing")
        else:
            correct_spacing_count = 0
            for text, _, line_spacing in all_question_paragraphs:
                if line_spacing is not None and abs(float(line_spacing) - TARGET_LINE_SPACING) < 0.01:
                    correct_spacing_count += 1
                else:
                    print(f"  Line spacing mismatch for '{text[:40]}...': expected {TARGET_LINE_SPACING}, got {line_spacing}")

            if correct_spacing_count == len(all_question_paragraphs) and correct_spacing_count > 0:
                print(f"PASS: Component 3 -- All {correct_spacing_count} questions have 2.0 line spacing (0.3 pts)")
                total_score += 0.3
            elif correct_spacing_count > 0:
                partial = round(0.3 * correct_spacing_count / len(all_question_paragraphs), 2)
                print(f"PARTIAL: Component 3 -- {correct_spacing_count}/{len(all_question_paragraphs)} correct line spacing (partial {partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 3 -- No questions have correct 2.0 line spacing")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
