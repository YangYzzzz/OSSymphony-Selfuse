"""
Initial Setup: Ecology Seminar presentation with 10 slides, slide 9 empty with title only
Task ID: impress_teach_072
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_072'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Ecology Seminar: Understanding Our Living Planet"
    slide1.placeholders[1].text = "Department of Environmental Science\nSpring 2025 Lecture Series"

    # --- Slide 2: Introduction ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "What Is Ecology?"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Ecology is the scientific study of the interactions between organisms and their environment."
    p2 = tf2.add_paragraph()
    p2.text = "It encompasses the distribution and abundance of living organisms across ecosystems."
    p3 = tf2.add_paragraph()
    p3.text = "Key branches include population ecology, community ecology, and ecosystem ecology."
    p4 = tf2.add_paragraph()
    p4.text = "Understanding ecology helps us predict and manage environmental changes."

    # --- Slide 3: Biomes Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Major Terrestrial Biomes"
    tf3 = slide3.placeholders[1].text_frame
    biomes = [
        "Tropical Rainforest: High biodiversity, 2000+ mm rainfall annually",
        "Temperate Deciduous Forest: Seasonal leaf drop, moderate climate",
        "Boreal Forest (Taiga): Coniferous trees, long cold winters",
        "Grasslands: Dominated by grasses, limited tree growth",
        "Desert: Less than 250 mm annual precipitation, extreme temperatures",
        "Tundra: Permafrost layer, minimal vegetation, Arctic regions",
    ]
    tf3.text = biomes[0]
    for b in biomes[1:]:
        p = tf3.add_paragraph()
        p.text = b

    # --- Slide 4: Biodiversity ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Biodiversity and Its Importance"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Genetic diversity: Variation within species allows adaptation to changing conditions"
    items4 = [
        "Species diversity: The variety of species within a habitat or region",
        "Ecosystem diversity: Range of ecosystems across a geographical area",
        "An estimated 8.7 million species exist on Earth; only 1.2 million catalogued",
        "Biodiversity provides ecosystem services worth $125-145 trillion per year globally",
    ]
    for item in items4:
        p = tf4.add_paragraph()
        p.text = item

    # --- Slide 5: Climate Change Impact ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Climate Change and Ecosystems"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Global temperatures have risen 1.1\u00b0C since pre-industrial levels"
    items5 = [
        "Arctic sea ice has declined by 13% per decade since 1979",
        "Sea levels are rising at 3.6 mm per year, threatening coastal habitats",
        "Shifts in species ranges: Many organisms moving poleward or to higher elevations",
        "Phenological mismatches disrupt pollination and predator-prey relationships",
    ]
    for item in items5:
        p = tf5.add_paragraph()
        p.text = item

    # --- Slide 6: Coral Reef Ecosystems ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Coral Reef Ecosystems Under Threat"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Coral reefs support 25% of all marine species despite covering less than 1% of the ocean floor"
    items6 = [
        "Ocean acidification reduces coral calcification rates by up to 40%",
        "Mass bleaching events have increased fivefold since the 1980s",
        "The Great Barrier Reef has lost over 50% of its coral cover since 1995",
        "Reef restoration efforts show promise but cannot replace natural recovery",
    ]
    for item in items6:
        p = tf6.add_paragraph()
        p.text = item

    # --- Slide 7: Human Impact ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Human Activities and Ecosystem Disruption"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "Deforestation: 10 million hectares of forest lost annually"
    items7 = [
        "Urbanization: 55% of the world population lives in cities, fragmenting habitats",
        "Pollution: 8 million tons of plastic enter oceans each year",
        "Overexploitation: 34% of global fish stocks are overfished",
        "Invasive species: Cost the global economy $423 billion per year",
    ]
    for item in items7:
        p = tf7.add_paragraph()
        p.text = item

    # --- Slide 8: Conservation Strategies ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Conservation Strategies"
    tf8 = slide8.placeholders[1].text_frame
    tf8.text = "Protected areas now cover 17% of terrestrial and 8% of marine environments"
    items8 = [
        "Wildlife corridors connect fragmented habitats, enabling gene flow",
        "Community-based conservation empowers local stewardship of resources",
        "Rewilding projects have successfully reintroduced wolves, beavers, and bison",
        "Payment for ecosystem services incentivizes landowners to preserve habitats",
    ]
    for item in items8:
        p = tf8.add_paragraph()
        p.text = item

    # --- Slide 9: Discussion Time (EMPTY - just title, no content) ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide9.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
    tf9 = txBox.text_frame
    p9 = tf9.paragraphs[0]
    p9.text = "Discussion Time"
    p9.alignment = PP_ALIGN.CENTER
    run9 = p9.runs[0]
    run9.font.size = Pt(36)
    run9.font.bold = True
    run9.font.color.rgb = RGBColor(0x2E, 0x4A, 0x62)

    # --- Slide 10: References & Further Reading ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "References & Further Reading"
    tf10 = slide10.placeholders[1].text_frame
    tf10.text = "IPCC (2023). Climate Change 2023: Synthesis Report."
    items10 = [
        "Wilson, E.O. (2016). Half-Earth: Our Planet's Fight for Life. Norton.",
        "Kolbert, E. (2014). The Sixth Extinction: An Unnatural History. Holt.",
        "IUCN Red List of Threatened Species: www.iucnredlist.org",
        "UN Environment Programme: www.unep.org/resources",
    ]
    for item in items10:
        p = tf10.add_paragraph()
        p.text = item

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
