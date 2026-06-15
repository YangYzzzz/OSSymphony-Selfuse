"""
Initial Setup: Create a 5-slide Skills Assessment presentation with no chart on slide 3.
Task ID: impress_tct_067
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_067'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=None, alignment=None):
    """Helper to add a text box to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Skills Assessment Report"
    slide1.placeholders[1].text = "Q1 2026 Team Evaluation"

    # --- Slide 2: Assessment Methodology ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Assessment Methodology"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Our team skills assessment uses a comprehensive 5-axis evaluation framework:"
    items = [
        "Technical: Coding proficiency, system design, debugging expertise",
        "Communication: Written clarity, presentation skills, active listening",
        "Leadership: Decision making, mentoring, strategic vision",
        "Creativity: Innovation, lateral thinking, novel problem approaches",
        "Problem Solving: Analytical reasoning, root cause analysis, solution design",
    ]
    for item in items:
        p = body2.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 3: Team Skills Comparison (NO chart — just title and data table) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title manually
    add_textbox(slide3, 0.5, 0.3, 9, 1, "Team Skills Comparison",
                font_size=28, bold=True, color=RGBColor(0x1F, 0x49, 0x7D),
                alignment=PP_ALIGN.CENTER)

    # Add a data table showing the raw scores
    tbl_shape = slide3.shapes.add_table(3, 6, Inches(0.8), Inches(1.8), Inches(8.4), Inches(1.5))
    tbl = tbl_shape.table

    headers = ["", "Technical", "Communication", "Leadership", "Creativity", "Problem Solving"]
    data_rows = [
        ["Alex Rivera", "9", "7", "6", "8", "9"],
        ["Jordan Kim", "7", "9", "8", "7", "6"],
    ]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(12)
    for r, row_data in enumerate(data_rows, 1):
        for c, val in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(12)

    # Add a note below the table
    add_textbox(slide3, 0.8, 3.8, 8.4, 0.8,
                "Scores are on a 1-10 scale based on peer review, self-assessment, and manager evaluation.",
                font_size=11, color=RGBColor(0x66, 0x66, 0x66))

    # --- Slide 4: Individual Development Plans ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Individual Development Plans"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Alex Rivera - Focus Areas:"
    plans_alex = [
        "Leadership workshop enrollment (Q2 2026)",
        "Cross-functional project lead assignment",
        "Communication skills coaching sessions",
    ]
    for item in plans_alex:
        p = body4.add_paragraph()
        p.text = item
        p.level = 1

    p_jordan = body4.add_paragraph()
    p_jordan.text = "Jordan Kim - Focus Areas:"
    plans_jordan = [
        "Advanced technical certification program",
        "Problem-solving methodology training",
        "Innovation hackathon participation",
    ]
    for item in plans_jordan:
        p = body4.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 5: Next Steps & Timeline ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Next Steps & Timeline"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Quarterly Review Schedule:"
    steps = [
        "April 2026: Development plan kickoff meetings",
        "June 2026: Mid-quarter progress check-ins",
        "July 2026: Q2 skills reassessment",
        "September 2026: Annual performance review integration",
        "December 2026: Year-end comprehensive evaluation",
    ]
    for item in steps:
        p = body5.add_paragraph()
        p.text = item
        p.level = 1

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
