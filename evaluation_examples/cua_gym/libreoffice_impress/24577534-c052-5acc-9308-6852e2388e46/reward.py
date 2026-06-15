"""
Reward Script: Create radar/spider chart on slide 3 comparing skill ratings
Task ID: impress_tct_067
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Slide 3 contains a radar-type chart
  Component 2 (0.20): Chart has 5 correct skill categories
  Component 3 (0.20): Chart has 2 series with correct data values
  Component 4 (0.15): Series names reference the two team members
  Component 5 (0.15): Chart has a legend enabled
"""

import os
import sys

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_067'

# Expected skill categories (order matters for radar charts)
EXPECTED_CATEGORIES = ['Technical', 'Communication', 'Leadership', 'Creativity', 'Problem Solving']

# Expected series data (from the table on initial slide 3)
EXPECTED_SERIES = {
    'Alex Rivera': [9.0, 7.0, 6.0, 8.0, 9.0],
    'Jordan Kim': [7.0, 9.0, 8.0, 7.0, 6.0],
}

EXPECTED_MEMBER_NAMES = ['Alex Rivera', 'Jordan Kim']


def persist_app_state():
    """Try to save any unsaved LibreOffice state."""
    try:
        os.environ["DISPLAY"] = ":0"
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        import time
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.chart import XL_CHART_TYPE
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"PRECONDITION FAIL: Need at least 3 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]

    # Find chart shape(s) on slide 3
    chart_shape = None
    for shape in slide3.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Slide 3 contains a radar-type chart (0.30 points)
    # Radar chart types: RADAR (81), RADAR_FILLED (82), RADAR_MARKERS (83)
    RADAR_TYPES = {81, 82, 83}
    try:
        if chart_shape is not None:
            chart_type_val = chart_shape.chart.chart_type.real
            if chart_type_val in RADAR_TYPES:
                print(f"PASS: Component 1 - Radar chart found on slide 3 (type={chart_type_val}) (0.30 pts)")
                total_score += 0.30
            else:
                print(f"FAIL: Component 1 - Chart found but not radar type (type={chart_type_val})")
        else:
            print("FAIL: Component 1 - No chart found on slide 3")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # If no chart found, remaining components cannot be checked
    if chart_shape is None:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    chart = chart_shape.chart

    # Component 2: Chart has 5 correct skill categories (0.20 points)
    try:
        plot = chart.plots[0]
        categories = list(plot.categories)
        if len(categories) == 5:
            # Check category names match (case-insensitive, stripped)
            expected_lower = [c.strip().lower() for c in EXPECTED_CATEGORIES]
            actual_lower = [c.strip().lower() for c in categories]
            if actual_lower == expected_lower:
                print(f"PASS: Component 2 - All 5 categories match: {categories} (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 2 - Categories don't match. Expected: {EXPECTED_CATEGORIES}, Got: {categories}")
        else:
            print(f"FAIL: Component 2 - Expected 5 categories, found {len(categories)}: {categories}")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Chart has 2 series with correct data values (0.20 points)
    try:
        series_list = list(chart.series)
        if len(series_list) == 2:
            series0_vals = list(series_list[0].values)
            series1_vals = list(series_list[1].values)

            # Check both orderings (series could be in either order)
            expected_vals_list = list(EXPECTED_SERIES.values())
            match_order1 = (series0_vals == expected_vals_list[0] and series1_vals == expected_vals_list[1])
            match_order2 = (series0_vals == expected_vals_list[1] and series1_vals == expected_vals_list[0])

            if match_order1 or match_order2:
                print(f"PASS: Component 3 - 2 series with correct values (0.20 pts)")
                print(f"  Series 0: {series0_vals}")
                print(f"  Series 1: {series1_vals}")
                total_score += 0.20
            else:
                print(f"FAIL: Component 3 - Series values don't match expected")
                print(f"  Expected: {expected_vals_list}")
                print(f"  Got: [{series0_vals}, {series1_vals}]")
        else:
            print(f"FAIL: Component 3 - Expected 2 series, found {len(series_list)}")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Series names reference the two team members (0.15 points)
    try:
        # Extract series names from chart XML
        chart_xml = chart.part._element.xml
        import xml.etree.ElementTree as ET
        ns = {
            'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
        }
        root = ET.fromstring(chart_xml)
        series_names = []
        for ser in root.findall('.//c:ser', ns):
            tx = ser.find('c:tx', ns)
            if tx is not None:
                str_cache = tx.find('.//c:strCache', ns)
                if str_cache is not None:
                    pt = str_cache.find('.//c:pt/c:v', ns)
                    if pt is not None and pt.text:
                        series_names.append(pt.text.strip())

        # Check if both team member names appear
        found_names = set()
        for name in series_names:
            name_lower = name.lower()
            for member in EXPECTED_MEMBER_NAMES:
                if member.lower() in name_lower or name_lower in member.lower():
                    found_names.add(member)

        if len(found_names) == 2:
            print(f"PASS: Component 4 - Both team members found in series names: {series_names} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 - Expected both team members in series names. Found: {series_names}, Matched: {found_names}")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Chart has a legend enabled (0.15 points)
    try:
        if chart.has_legend:
            print(f"PASS: Component 5 - Chart legend is enabled (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 5 - Chart has no legend")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
