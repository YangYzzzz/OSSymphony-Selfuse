"""
FINAL REWARD SCRIPT - SUCCESS
Task: Set Table 1 column widths to 3.0 cm, 2.0 cm, 4.0 cm for the first three columns.
Generated: 2025-10-17 11:16:31
Status: success
Model: azure-o3
Total Steps: 1
"""

from pptx import Presentation
import os

# Helper: Convert English Metric Units (EMU) to centimeters
def emu_to_cm(emu):
    return emu / 360000.0

def verify_table_column_widths(file_path):
    """Verify that the first table in a PPTX has its first three columns
    set to 3.0 cm, 2.0 cm, and 4.0 cm respectively.
    Progressive scoring is applied:
      • 0.25 pts – table with ≥3 columns exists
      • 0.25 pts – column 1 within tolerance
      • 0.25 pts – column 2 within tolerance
      • 0.25 pts – column 3 within tolerance
    Returns a score between 0.0 and 1.0 and prints detailed diagnostics.
    """

    max_score = 1.0
    total_score = 0.0
    tolerance = 0.15  # ±0.15 cm considered acceptable

    # 1. Check file presence
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    # 2. Load presentation
    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded presentation with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # 3. Locate first table (Table 1)
    table = None
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if shape.has_table:
                table = shape.table
                print(f"✓ Found table on slide {slide_idx} shape {shape_idx}")
                break
        if table is not None:
            break

    if table is None:
        print("✗ No table found in the presentation")
        print("REWARD: 0.0")
        return 0.0

    # 4. Verify column count
    if len(table.columns) >= 3:
        total_score += 0.25
        print("✓ Table has at least 3 columns (+0.25)")
    else:
        print("✗ Table doesn’t have enough columns – cannot verify widths")
        print(f"REWARD: {total_score}")
        return total_score

    # 5. Expected widths (cm)
    expected = [3.0, 2.0, 4.0]
    column_weights = 0.25  # per-column weight

    for idx in range(3):
        actual_cm = emu_to_cm(table.columns[idx].width)
        diff = abs(actual_cm - expected[idx])
        print(f"Column {idx+1}: {actual_cm:.2f} cm (expected {expected[idx]} cm) diff {diff:.2f} cm")
        if diff <= tolerance:
            total_score += column_weights
            print(f"✓ Column {idx+1} within tolerance (+{column_weights})")
        else:
            print(f"✗ Column {idx+1} out of tolerance (0 pts)")

    final_score = min(total_score, max_score)
    print(f"Total score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score

# ---------- Execute verification ----------
file_path = '/home/user/set_table_1_column_widths_to_30_cm_20_cm_40_cm_for_the_first_three_columns.pptx'
verify_table_column_widths(file_path)
