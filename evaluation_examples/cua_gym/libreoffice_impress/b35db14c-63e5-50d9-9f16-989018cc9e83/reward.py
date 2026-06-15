"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 4, how do I draw a rectangle that’s exactly 6 cm wide and 2 cm high, fill it with the preset color “Green 6” (hex #00A933), and make sure it has no outline at all?
Generated: 2025-09-10 14:21:24
Status: success
Model: azure-o3
Total Steps: 6
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL_TYPE


def verify_rectangle_task(file_path: str) -> float:
    """Verify that on slide 4 there is a 6 cm × 2 cm rectangle
    filled with Green 6 (#00A933) and with no outline.

    Returns a progressive score from 0.0–1.0 and prints details.
    """

    MAX_SCORE = 1.0
    score = 0.0

    # Conversion and tolerance
    CM_TO_EMU = 360_000          # 1 cm in English Metric Units
    TOL_EMU = 10_000             # ≈0.28 mm tolerance
    req_w = 6 * CM_TO_EMU        # 6 cm in EMU
    req_h = 2 * CM_TO_EMU        # 2 cm in EMU

    # ------------------------------------------------------------------
    # 0. Load file ------------------------------------------------------
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Could not open presentation: {e}")
        return 0.0

    # ------------------------------------------------------------------
    # 1. Ensure slide 4 exists -----------------------------------------
    # ------------------------------------------------------------------
    if len(prs.slides) < 4:
        print(f"✗ Need at least 4 slides, found {len(prs.slides)}")
        return 0.0

    slide = prs.slides[3]  # zero-based index → slide 4

    # ------------------------------------------------------------------
    # 2. Locate rectangle with correct dimensions ----------------------
    # ------------------------------------------------------------------
    rectangle = None
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue  # not an auto-shape, skip
        if abs(shape.width - req_w) <= TOL_EMU and abs(shape.height - req_h) <= TOL_EMU:
            rectangle = shape
            break

    if rectangle is None:
        print("✗ Rectangle 6 cm × 2 cm not found on slide 4")
        return 0.0

    print("✓ Found rectangle with correct dimensions (0.4 points)")
    score += 0.4

    # ------------------------------------------------------------------
    # 3. Verify fill colour --------------------------------------------
    # ------------------------------------------------------------------
    fill_ok = False
    fill = rectangle.fill
    if fill.type is not None:
        try:
            rgb = str(fill.fore_color.rgb).upper()  # returns e.g. '00A933'
            if rgb == "00A933":
                fill_ok = True
        except Exception:
            pass

    if fill_ok:
        print("✓ Fill colour is Green 6 (#00A933) (0.3 points)")
        score += 0.3
    else:
        print("✗ Fill colour incorrect or not set to #00A933")

    # ------------------------------------------------------------------
    # 4. Verify absence of outline -------------------------------------
    # ------------------------------------------------------------------
    no_outline = False
    try:
        line = rectangle.line
        # width 0 or fill type NONE/BACKGROUND/None all indicate no visible outline
        if line.width == 0:
            no_outline = True
        else:
            ft = line.fill.type  # can be None, BACKGROUND, etc.
            if ft is None or ft == MSO_FILL_TYPE.BACKGROUND:
                no_outline = True
    except Exception as e:
        print("  Warning while checking outline:", e)

    if no_outline:
        print("✓ No outline detected (0.3 points)")
        score += 0.3
    else:
        print("✗ Outline present")

    # ------------------------------------------------------------------
    # 5. Final score ----------------------------------------------------
    # ------------------------------------------------------------------
    final_score = min(score, MAX_SCORE)
    print(f"Total score: {final_score}")
    return final_score


if __name__ == "__main__":
    pptx_path = (
        "/home/user/on_slide_4_how_do_i_draw_a_rectangle_thats_exactly_6_cm_" \
        "wide_and_2_cm_high_fill_it_with_the_preset_c_golden.pptx"
    )
    reward = verify_rectangle_task(pptx_path)
    print(f"REWARD: {reward}")
