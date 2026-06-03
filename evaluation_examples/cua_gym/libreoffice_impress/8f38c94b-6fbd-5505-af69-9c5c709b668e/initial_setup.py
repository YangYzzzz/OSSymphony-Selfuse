"""
Initial Setup: Create a 6-slide Medical Case presentation with empty Slide 4
Task ID: impress_teach_065
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_065'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Comprehensive Medical Case Study"
    slide1.placeholders[1].text = "Department of Internal Medicine\nSt. Margaret's Hospital\nPresented by Dr. Elena Vasquez"

    # --- Slide 2: Patient History ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide2.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Patient History"
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    content_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(5))
    tf2 = content_box.text_frame
    tf2.word_wrap = True
    items = [
        "Patient: Maria Gonzalez, 54-year-old female",
        "Chief Complaint: Progressive joint pain and morning stiffness for 6 months",
        "Medical History: Type 2 diabetes (diagnosed 2018), hypertension",
        "Family History: Mother had rheumatoid arthritis",
        "Medications: Metformin 1000mg BID, Lisinopril 10mg daily",
        "Allergies: Penicillin (rash), Sulfa drugs",
        "Social History: Non-smoker, occasional alcohol use, retired teacher",
        "Review of Systems: Fatigue, weight loss of 8 lbs over 3 months",
    ]
    for i, item in enumerate(items):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = item
        p.space_after = Pt(8)
        for r in p.runs:
            r.font.size = Pt(16)

    # --- Slide 3: Diagnosis ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Diagnosis & Assessment"
    run3 = p3.runs[0]
    run3.font.size = Pt(32)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    diag_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(5))
    tf3b = diag_box.text_frame
    tf3b.word_wrap = True
    diag_items = [
        "Primary Diagnosis: Rheumatoid Arthritis (ICD-10: M06.9)",
        "RF Factor: Positive (128 IU/mL, reference <14)",
        "Anti-CCP Antibodies: Strongly positive (>250 U/mL)",
        "ESR: 48 mm/hr (elevated), CRP: 3.2 mg/dL (elevated)",
        "X-Ray findings: Bilateral MCP joint erosions, periarticular osteopenia",
        "DAS28 Score: 5.8 (High disease activity)",
        "Comorbidity considerations: Diabetes management during immunosuppression",
    ]
    for i, item in enumerate(diag_items):
        if i == 0:
            p = tf3b.paragraphs[0]
        else:
            p = tf3b.add_paragraph()
        p.text = item
        p.space_after = Pt(8)
        for r in p.runs:
            r.font.size = Pt(16)

    # --- Slide 4: Treatment Results (EMPTY - just title) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title_box = slide4.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(10), Inches(0.8))
    tf4 = title_box.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Treatment Results"
    run4 = p4.runs[0]
    run4.font.size = Pt(32)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    # --- Slide 5: Follow-Up Plan ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox5 = slide5.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
    tf5 = txBox5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Follow-Up Plan"
    run5 = p5.runs[0]
    run5.font.size = Pt(32)
    run5.font.bold = True
    run5.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    plan_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(5))
    tf5b = plan_box.text_frame
    tf5b.word_wrap = True
    plan_items = [
        "Week 2: Phone follow-up - assess methotrexate tolerance and side effects",
        "Week 4: Lab work - CBC, LFTs, renal function panel",
        "Week 8: In-office visit - DAS28 reassessment, adjust therapy if needed",
        "Month 3: Repeat imaging of hands and feet",
        "Month 6: Comprehensive disease activity evaluation",
        "Ongoing: Monthly CBC and LFTs for first 6 months of MTX therapy",
        "Referral to occupational therapy for joint protection strategies",
        "Diabetes monitoring: HbA1c every 3 months during steroid taper",
    ]
    for i, item in enumerate(plan_items):
        if i == 0:
            p = tf5b.paragraphs[0]
        else:
            p = tf5b.add_paragraph()
        p.text = item
        p.space_after = Pt(8)
        for r in p.runs:
            r.font.size = Pt(16)

    # --- Slide 6: References ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox6 = slide6.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
    tf6 = txBox6.text_frame
    p6 = tf6.paragraphs[0]
    p6.text = "References"
    run6 = p6.runs[0]
    run6.font.size = Pt(32)
    run6.font.bold = True
    run6.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

    ref_box = slide6.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(5))
    tf6b = ref_box.text_frame
    tf6b.word_wrap = True
    refs = [
        "1. Aletaha D, et al. 2010 Rheumatoid arthritis classification criteria. Arthritis Rheum. 2010;62(9):2569-2581.",
        "2. Singh JA, et al. 2015 ACR Guideline for RA Treatment. Arthritis Rheumatol. 2016;68(1):1-26.",
        "3. Smolen JS, et al. EULAR recommendations for RA management. Ann Rheum Dis. 2023;82:3-18.",
        "4. Burmester GR, Pope JE. Novel treatment strategies in RA. Lancet. 2017;389(10086):2338-2348.",
        "5. Fraenkel L, et al. 2021 ACR Guideline for RA Treatment. Arthritis Care Res. 2021;73(7):924-939.",
    ]
    for i, item in enumerate(refs):
        if i == 0:
            p = tf6b.paragraphs[0]
        else:
            p = tf6b.add_paragraph()
        p.text = item
        p.space_after = Pt(6)
        for r in p.runs:
            r.font.size = Pt(14)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
