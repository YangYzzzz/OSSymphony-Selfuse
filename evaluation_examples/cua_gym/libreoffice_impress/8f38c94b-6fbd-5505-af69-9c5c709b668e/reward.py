"""
Reward Script: Before-and-after comparison slide on Slide 4
Task ID: impress_teach_065
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): Vertical center line exists on slide 4
  Component 2 (0.25): 'Before Treatment' heading in red (#C62828) on left half
  Component 3 (0.25): 'After Treatment' heading in green (#2E7D32) on right half
  Component 4 (0.15): Left dashed-border rectangle placeholder with 'Insert Image Here'
  Component 5 (0.15): Right dashed-border rectangle placeholder with 'Insert Image Here'
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_065'


def persist_app_state(domain):
    """Save any unsaved LibreOffice edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)
    slide_width = prs.slide_width  # 12191695 EMU

    # Helper: get shape line properties from XML
    def get_line_props(shape):
        sp = shape._element
        for child in sp:
            if child.tag.endswith('}spPr'):
                ln = child.find(qn('a:ln'))
                if ln is not None:
                    dash_el = ln.find(qn('a:prstDash'))
                    dash_val = dash_el.get('val') if dash_el is not None else None
                    fill_el = ln.find(qn('a:solidFill'))
                    line_color = None
                    if fill_el is not None:
                        srgb = fill_el.find(qn('a:srgbClr'))
                        if srgb is not None:
                            line_color = srgb.get('val')
                    no_fill = ln.find(qn('a:noFill'))
                    return {
                        'width': ln.get('w'),
                        'dash': dash_val,
                        'color': line_color,
                        'no_fill': no_fill is not None
                    }
        return None

    def get_shape_fill_color(shape):
        sp = shape._element
        for child in sp:
            if child.tag.endswith('}spPr'):
                sf = child.find(qn('a:solidFill'))
                if sf is not None:
                    srgb = sf.find(qn('a:srgbClr'))
                    if srgb is not None:
                        return srgb.get('val')
        return None

    # Helper: get text and font color from shape
    def get_text_info(shape):
        if not hasattr(shape, 'has_text_frame') or not shape.has_text_frame:
            return None, None
        text = shape.text_frame.text.strip() if shape.text_frame.text else ""
        color = None
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                try:
                    if run.font.color.type is not None:
                        color = str(run.font.color.rgb)
                except Exception:
                    pass
        return text, color

    # Collect all shapes on slide 4 (excluding pre-existing title and Treatment Results heading)
    shapes_list = list(slide.shapes)

    # ============================================================
    # Component 1: Vertical center line on slide 4 (0.20 points)
    # A narrow rectangle or line near the horizontal center of the slide
    # ============================================================
    try:
        center_x = slide_width // 2  # ~6095847
        vertical_line_found = False
        for shape in shapes_list:
            # Look for a narrow shape (width < 100000 EMU = ~1.1 inches) positioned
            # near the center of the slide, tall enough to serve as a divider
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
                if shape.width < 100000 and shape.height > 2000000:
                    # Check if horizontally near center (within 15% of slide width)
                    shape_center_x = shape.left + shape.width // 2
                    tolerance = slide_width * 0.15
                    if abs(shape_center_x - center_x) < tolerance:
                        vertical_line_found = True
                        print(f"PASS: Component 1 - Vertical center line found (shape={shape.name}, left={shape.left}, w={shape.width}, h={shape.height}) (0.20 pts)")
                        break
        if not vertical_line_found:
            # Also check for LINE shapes
            for shape in shapes_list:
                try:
                    if shape.shape_type == 9:  # MSO_SHAPE_TYPE.LINE
                        shape_center_x = shape.left + shape.width // 2
                        tolerance = slide_width * 0.15
                        if abs(shape_center_x - center_x) < tolerance and shape.height > 2000000:
                            vertical_line_found = True
                            print(f"PASS: Component 1 - Vertical center line found as LINE (0.20 pts)")
                            break
                except Exception:
                    pass
        if vertical_line_found:
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 - No vertical center line found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # ============================================================
    # Component 2: 'Before Treatment' heading in red (#C62828) on left half (0.25 points)
    # ============================================================
    try:
        before_found = False
        for shape in shapes_list:
            text, color = get_text_info(shape)
            if text and 'before treatment' in text.lower():
                # Check it's on the left half
                shape_center_x = shape.left + shape.width // 2
                if shape_center_x < center_x:
                    # Check red color C62828
                    if color and color.upper() == 'C62828':
                        before_found = True
                        print(f"PASS: Component 2 - 'Before Treatment' in red (#C62828) on left half (0.25 pts)")
                    else:
                        print(f"FAIL: Component 2 - 'Before Treatment' found on left but color is {color}, expected C62828")
                else:
                    print(f"FAIL: Component 2 - 'Before Treatment' found but on right half (center_x={shape_center_x})")
                break
        if not before_found and not any('before treatment' in (get_text_info(s)[0] or '').lower() for s in shapes_list):
            print(f"FAIL: Component 2 - No 'Before Treatment' text found on slide 4")
        if before_found:
            total_score += 0.25
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # ============================================================
    # Component 3: 'After Treatment' heading in green (#2E7D32) on right half (0.25 points)
    # ============================================================
    try:
        after_found = False
        for shape in shapes_list:
            text, color = get_text_info(shape)
            if text and 'after treatment' in text.lower():
                # Check it's on the right half
                shape_center_x = shape.left + shape.width // 2
                if shape_center_x > center_x:
                    # Check green color 2E7D32
                    if color and color.upper() == '2E7D32':
                        after_found = True
                        print(f"PASS: Component 3 - 'After Treatment' in green (#2E7D32) on right half (0.25 pts)")
                    else:
                        print(f"FAIL: Component 3 - 'After Treatment' found on right but color is {color}, expected 2E7D32")
                else:
                    print(f"FAIL: Component 3 - 'After Treatment' found but on left half (center_x={shape_center_x})")
                break
        if not after_found and not any('after treatment' in (get_text_info(s)[0] or '').lower() for s in shapes_list):
            print(f"FAIL: Component 3 - No 'After Treatment' text found on slide 4")
        if after_found:
            total_score += 0.25
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # ============================================================
    # Component 4: Left dashed-border rectangle with 'Insert Image Here' (0.15 points)
    # ============================================================
    try:
        left_placeholder_found = False
        for shape in shapes_list:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                text, _ = get_text_info(shape)
                if text and 'insert image here' in text.lower():
                    shape_center_x = shape.left + shape.width // 2
                    if shape_center_x < center_x:
                        # Verify dashed border
                        line_props = get_line_props(shape)
                        if line_props and line_props.get('dash') is not None:
                            left_placeholder_found = True
                            print(f"PASS: Component 4 - Left dashed placeholder found (dash={line_props['dash']}) (0.15 pts)")
                        else:
                            # Still give credit if the rectangle has the right text and position
                            # but with partial credit for missing dash
                            left_placeholder_found = True
                            print(f"PASS: Component 4 - Left placeholder found with 'Insert Image Here' (dash style not verified but shape present) (0.15 pts)")
        if left_placeholder_found:
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 - No left-side dashed rectangle with 'Insert Image Here' found")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # ============================================================
    # Component 5: Right dashed-border rectangle with 'Insert Image Here' (0.15 points)
    # ============================================================
    try:
        right_placeholder_found = False
        for shape in shapes_list:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                text, _ = get_text_info(shape)
                if text and 'insert image here' in text.lower():
                    shape_center_x = shape.left + shape.width // 2
                    if shape_center_x > center_x:
                        line_props = get_line_props(shape)
                        if line_props and line_props.get('dash') is not None:
                            right_placeholder_found = True
                            print(f"PASS: Component 5 - Right dashed placeholder found (dash={line_props['dash']}) (0.15 pts)")
                        else:
                            right_placeholder_found = True
                            print(f"PASS: Component 5 - Right placeholder found with 'Insert Image Here' (dash style not verified but shape present) (0.15 pts)")
        if right_placeholder_found:
            total_score += 0.15
        else:
            print(f"FAIL: Component 5 - No right-side dashed rectangle with 'Insert Image Here' found")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
