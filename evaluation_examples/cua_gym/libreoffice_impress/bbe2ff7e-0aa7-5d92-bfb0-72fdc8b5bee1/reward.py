"""
Reward Script: Design a dark-theme corporate master slide
Task ID: impress_gf2_024
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25) — Master slide background is solid #111827
  Component 2 (0.25) — Title placeholder: 40pt bold white, Liberation Sans font
  Component 3 (0.20) — Content placeholder: 18pt, color #D1D5DB
  Component 4 (0.15) — Accent bar: 0.4cm wide, full height, #3B82F6, at left edge
  Component 5 (0.15) — Slide number placeholder: 12pt, color #6B7280
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_024'

# XML namespaces used in pptx
NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}


def get_master_xml(pptx_path):
    """Parse slideMaster1.xml from the pptx zip."""
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        with zf.open('ppt/slideMasters/slideMaster1.xml') as f:
            return ET.parse(f).getroot()


def find_placeholder_sp(root, ph_type):
    """Find the <p:sp> element for a given placeholder type (title, body, sldNum, etc.)."""
    for sp in root.findall('.//p:sp', NS):
        ph = sp.find('.//p:ph', NS)
        if ph is not None and ph.get('type') == ph_type:
            return sp
    return None


def get_solid_fill_color(element):
    """Extract srgbClr val from a solidFill child, if present."""
    solid = element.find('.//a:solidFill', NS)
    if solid is not None:
        srgb = solid.find('a:srgbClr', NS)
        if srgb is not None:
            return srgb.get('val', '').upper()
    return None


def verify_task(file_path):
    """
    Verify dark-theme corporate master slide with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        root = get_master_xml(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot parse master slide XML from {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Also load via python-pptx for higher-level checks
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        master = prs.slide_masters[0]
    except Exception as e:
        print(f"CRITICAL: Cannot load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # =========================================================================
    # Component 1: Master slide background is solid #111827 (0.25 points)
    # Initial: background fill type is BACKGROUND(5)/inherited, NOT solid #111827
    # Golden: solid fill #111827
    # =========================================================================
    try:
        bg_fill = master.background.fill
        if bg_fill.type is not None and bg_fill.type == 1:  # SOLID
            bg_color = str(bg_fill.fore_color.rgb).upper()
            if bg_color == '111827':
                print(f"PASS: Component 1 — Master background is solid #111827 (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 — Master background color is #{bg_color}, expected #111827")
        else:
            print(f"FAIL: Component 1 — Master background fill type is {bg_fill.type}, expected SOLID (1)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # =========================================================================
    # Component 2: Title placeholder — 40pt bold white, Liberation Sans (0.25 pts)
    # Initial: default font, no custom size/bold/color
    # Golden: Liberation Sans 40pt bold #FFFFFF
    # =========================================================================
    try:
        title_sp = find_placeholder_sp(root, 'title')
        if title_sp is None:
            print("FAIL: Component 2 — Title placeholder not found on master")
        else:
            c2_score = 0.0
            # Check via defRPr and rPr on the title placeholder
            # Look for size=4000 (40pt), b=1, color=FFFFFF
            all_defrpr = title_sp.findall('.//a:defRPr', NS)
            all_rpr = title_sp.findall('.//a:rPr', NS)
            rpr_elements = all_defrpr + all_rpr

            found_size_40 = False
            found_bold = False
            found_white = False
            found_font = False

            for rpr in rpr_elements:
                sz = rpr.get('sz')
                if sz and int(sz) == 4000:
                    found_size_40 = True
                b = rpr.get('b')
                if b == '1':
                    found_bold = True
                color = get_solid_fill_color(rpr)
                if color == 'FFFFFF':
                    found_white = True

            # Check font name via python-pptx
            title_shape = master.shapes[0]  # Title Placeholder 1
            if hasattr(title_shape, 'text_frame'):
                for para in title_shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name and run.font.name in ('Liberation Sans', 'Helvetica Neue'):
                            found_font = True

            # Also check font in XML (latin typeface)
            for latin in title_sp.findall('.//a:latin', NS):
                typeface = latin.get('typeface', '')
                if typeface in ('Liberation Sans', 'Helvetica Neue'):
                    found_font = True

            sub_checks = [found_size_40, found_bold, found_white, found_font]
            # Need all 4 for full points, partial credit for partial
            passed = sum(sub_checks)
            if passed == 4:
                c2_score = 0.25
            elif passed >= 2:
                c2_score = 0.15
            elif passed >= 1:
                c2_score = 0.05

            if c2_score > 0:
                print(f"PASS: Component 2 — Title: size40={found_size_40}, bold={found_bold}, white={found_white}, font={found_font} ({c2_score} pts)")
            else:
                print(f"FAIL: Component 2 — Title: size40={found_size_40}, bold={found_bold}, white={found_white}, font={found_font}")
            total_score += c2_score
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # =========================================================================
    # Component 3: Content/body placeholder — 18pt, color #D1D5DB (0.20 pts)
    # Initial: default sizes/colors (no custom defRPr color)
    # Golden: defRPr sz=1800 with solidFill #D1D5DB
    # =========================================================================
    try:
        body_sp = find_placeholder_sp(root, 'body')
        if body_sp is None:
            print("FAIL: Component 3 — Body placeholder not found on master")
        else:
            found_size_18 = False
            found_color_d1d5db = False

            for defrpr in body_sp.findall('.//a:defRPr', NS):
                sz = defrpr.get('sz')
                if sz and int(sz) == 1800:
                    found_size_18 = True
                color = get_solid_fill_color(defrpr)
                if color == 'D1D5DB':
                    found_color_d1d5db = True

            c3_score = 0.0
            if found_size_18 and found_color_d1d5db:
                c3_score = 0.20
            elif found_color_d1d5db:
                c3_score = 0.10
            elif found_size_18:
                c3_score = 0.05

            if c3_score > 0:
                print(f"PASS: Component 3 — Body: size18={found_size_18}, color_D1D5DB={found_color_d1d5db} ({c3_score} pts)")
            else:
                print(f"FAIL: Component 3 — Body: size18={found_size_18}, color_D1D5DB={found_color_d1d5db}")
            total_score += c3_score
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # =========================================================================
    # Component 4: Accent bar — 0.4cm wide, full height, #3B82F6, at left edge (0.15 pts)
    # Initial: no accent bar shape on master (only 5 shapes)
    # Golden: AccentBar shape at left=0, top=0, width=144000 EMU (0.4cm), height=6858000
    # =========================================================================
    try:
        # Look for a non-placeholder rectangle shape on the master
        accent_found = False
        accent_correct_size = False
        accent_correct_color = False
        accent_correct_pos = False

        cSld = root.find('.//p:cSld', NS)
        spTree = cSld.find('.//p:spTree', NS) if cSld is not None else None

        if spTree is not None:
            for sp in spTree.findall('p:sp', NS):
                # Skip placeholders
                ph = sp.find('.//p:ph', NS)
                if ph is not None:
                    continue

                # This is a non-placeholder shape — check if it's the accent bar
                xfrm = sp.find('.//a:xfrm', NS)
                if xfrm is None:
                    continue

                off = xfrm.find('a:off', NS)
                ext = xfrm.find('a:ext', NS)
                if off is None or ext is None:
                    continue

                x = int(off.get('x', '-1'))
                y = int(off.get('y', '-1'))
                cx = int(ext.get('cx', '0'))
                cy = int(ext.get('cy', '0'))

                # Check dimensions: width ~144000 EMU (0.4cm), height ~6858000 (full slide)
                # Use tolerance for width (within 10%)
                width_ok = abs(cx - 144000) <= 14400  # within 10%
                height_ok = abs(cy - 6858000) <= 68580  # within 1% of slide height
                pos_ok = x <= 5000 and y <= 5000  # near origin

                if width_ok and height_ok:
                    accent_found = True
                    accent_correct_size = True
                    accent_correct_pos = pos_ok

                    # Check fill color
                    solid_fill = sp.find('.//a:solidFill', NS)
                    if solid_fill is not None:
                        srgb = solid_fill.find('a:srgbClr', NS)
                        if srgb is not None and srgb.get('val', '').upper() == '3B82F6':
                            accent_correct_color = True

        c4_score = 0.0
        if accent_found and accent_correct_color and accent_correct_pos:
            c4_score = 0.15
        elif accent_found and accent_correct_color:
            c4_score = 0.10
        elif accent_found:
            c4_score = 0.05

        if c4_score > 0:
            print(f"PASS: Component 4 — Accent bar found: size={accent_correct_size}, color={accent_correct_color}, pos={accent_correct_pos} ({c4_score} pts)")
        else:
            print(f"FAIL: Component 4 — Accent bar: found={accent_found}, color={accent_correct_color}, pos={accent_correct_pos}")
        total_score += c4_score
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # =========================================================================
    # Component 5: Slide number placeholder — 12pt, color #6B7280 (0.15 pts)
    # Initial: default slide number, no custom color
    # Golden: rPr sz=1200 with solidFill #6B7280
    # =========================================================================
    try:
        sldnum_sp = find_placeholder_sp(root, 'sldNum')
        if sldnum_sp is None:
            print("FAIL: Component 5 — Slide number placeholder not found")
        else:
            found_size_12 = False
            found_color_6b7280 = False

            # Check rPr elements (run properties) and defRPr
            for rpr in sldnum_sp.findall('.//a:rPr', NS):
                sz = rpr.get('sz')
                if sz and int(sz) == 1200:
                    found_size_12 = True
                color = get_solid_fill_color(rpr)
                if color and color == '6B7280':
                    found_color_6b7280 = True

            for defrpr in sldnum_sp.findall('.//a:defRPr', NS):
                sz = defrpr.get('sz')
                if sz and int(sz) == 1200:
                    found_size_12 = True
                color = get_solid_fill_color(defrpr)
                if color and color == '6B7280':
                    found_color_6b7280 = True

            c5_score = 0.0
            # Note: initial already has defRPr sz=1200 but NO color.
            # So we require BOTH size AND color for any points (the color is the differentiator).
            if found_size_12 and found_color_6b7280:
                c5_score = 0.15
            elif found_color_6b7280:
                c5_score = 0.10

            if c5_score > 0:
                print(f"PASS: Component 5 — SlideNum: size12={found_size_12}, color_6B7280={found_color_6b7280} ({c5_score} pts)")
            else:
                print(f"FAIL: Component 5 — SlideNum: size12={found_size_12}, color_6B7280={found_color_6b7280}")
            total_score += c5_score
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entrypoint
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
