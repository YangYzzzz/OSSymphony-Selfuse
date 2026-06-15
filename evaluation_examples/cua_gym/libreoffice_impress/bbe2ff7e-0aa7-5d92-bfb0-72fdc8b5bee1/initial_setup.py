"""
Initial Setup: Dark-theme corporate presentation master slide design
Task ID: impress_gf2_024
Domain: libreoffice_impress

Creates a 20-slide presentation with a plain white master slide,
default black fonts, and no decorative elements.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_024'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard 16:9 slide dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Corporate slide content for 20 slides
    slide_content = [
        ("Q1 2025 Business Review", "Global Operations Division"),
        ("Agenda", "Financial overview\nMarket expansion\nTeam updates\nQ2 roadmap"),
        ("Revenue Summary", "Total revenue: $4.2M\nGrowth: 18% YoY\nNew clients: 47"),
        ("Regional Performance", "North America: $2.1M (+22%)\nEurope: $1.3M (+15%)\nAsia Pacific: $0.8M (+12%)"),
        ("Key Metrics Dashboard", "Customer retention: 94.3%\nNPS score: 72\nAverage deal size: $89K"),
        ("Product Roadmap", "Phase 1: Platform migration (Complete)\nPhase 2: AI integration (In progress)\nPhase 3: Mobile launch (Q3 2025)"),
        ("Team Highlights", "Engineering: 45 engineers, 3 new hires\nSales: Exceeded target by 12%\nSupport: Response time < 2 hours"),
        ("Customer Success Stories", "Meridian Corp: 40% efficiency gain\nAtlas Industries: $1.2M cost savings\nNova Tech: 3x user adoption"),
        ("Competitive Landscape", "Market share: 23% (up from 19%)\nKey differentiators: AI-powered analytics\nNew entrants: 2 in Q1"),
        ("Financial Projections", "Q2 target: $4.8M\nFull year forecast: $19.5M\nCapEx budget: $2.1M"),
        ("Marketing Initiatives", "Campaign reach: 2.3M impressions\nConversion rate: 3.8%\nBrand awareness: +15 points"),
        ("Technology Infrastructure", "Cloud migration: 85% complete\nUptime: 99.97%\nSecurity audits: All passed"),
        ("Partnership Updates", "Strategic partners: 12 active\nNew partnerships: Vertex AI, DataStream\nJoint revenue: $890K"),
        ("Risk Assessment", "Supply chain: Medium risk\nRegulatory: Low risk\nCybersecurity: Ongoing monitoring"),
        ("Employee Engagement", "Satisfaction score: 4.2/5.0\nRetention rate: 91%\nTraining hours: 2,400 total"),
        ("Sustainability Report", "Carbon offset: 150 tons\nPaperless initiative: 78% adoption\nEnergy efficiency: +25%"),
        ("Innovation Lab", "Patents filed: 3\nPrototypes: 5 in testing\nR&D investment: $1.5M"),
        ("Client Pipeline", "Qualified leads: 127\nProposal stage: 34\nExpected close rate: 42%"),
        ("Operational Efficiency", "Process automation: 60% of workflows\nCost per transaction: -15%\nSLA compliance: 98.5%"),
        ("Next Steps & Action Items", "Finalize Q2 budgets by April 15\nLaunch beta program May 1\nBoard presentation scheduled June 12"),
    ]

    for i, (title, body) in enumerate(slide_content):
        if i == 0:
            # Title slide layout
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title
            slide.placeholders[1].text = body
        else:
            # Title + Content layout
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title
            slide.placeholders[1].text = body

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Number of slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
