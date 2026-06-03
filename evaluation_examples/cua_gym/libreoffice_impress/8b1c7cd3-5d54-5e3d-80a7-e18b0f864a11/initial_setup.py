"""
Initial Setup: Insert product screenshot on slide 4 with border
Task ID: impress_sales_028
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_028'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
SCREENSHOT_PATH = f'{WORKDIR}/Desktop/dashboard_screenshot.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_screenshot():
    """Create a realistic 1920x1080 dashboard screenshot image."""
    os.makedirs(f'{WORKDIR}/Desktop', exist_ok=True)
    img = Image.new('RGB', (1920, 1080), color=(245, 247, 250))

    # Draw some rectangles to simulate a dashboard
    from PIL import ImageDraw, ImageFont
    draw = ImageDraw.Draw(img)

    # Header bar
    draw.rectangle([0, 0, 1920, 60], fill=(33, 37, 41))
    # Sidebar
    draw.rectangle([0, 60, 240, 1080], fill=(52, 58, 64))
    # Main content cards
    card_colors = [(255, 255, 255), (255, 255, 255), (255, 255, 255), (255, 255, 255)]
    positions = [(260, 80, 680, 280), (700, 80, 1120, 280), (1140, 80, 1560, 280), (1580, 80, 1900, 280)]
    for pos, color in zip(positions, card_colors):
        draw.rectangle(pos, fill=color, outline=(222, 226, 230), width=1)

    # Chart area
    draw.rectangle([260, 300, 1180, 700], fill=(255, 255, 255), outline=(222, 226, 230), width=1)
    # Bar chart simulation
    bar_colors = [(0, 123, 255), (40, 167, 69), (255, 193, 7), (220, 53, 69)]
    bar_x = 320
    for i in range(12):
        bar_height = 50 + (i * 17 + 30) % 250
        color = bar_colors[i % 4]
        draw.rectangle([bar_x, 650 - bar_height, bar_x + 55, 650], fill=color)
        bar_x += 70

    # Table area
    draw.rectangle([260, 720, 1180, 1050], fill=(255, 255, 255), outline=(222, 226, 230), width=1)
    for row_y in range(740, 1040, 30):
        draw.line([(280, row_y), (1160, row_y)], fill=(233, 236, 239), width=1)

    # Right panel
    draw.rectangle([1200, 300, 1900, 700], fill=(255, 255, 255), outline=(222, 226, 230), width=1)
    # Pie chart simulation (just a filled circle)
    draw.ellipse([1350, 350, 1750, 650], fill=(0, 123, 255), outline=(255, 255, 255))
    draw.pieslice([1350, 350, 1750, 650], start=0, end=120, fill=(40, 167, 69))
    draw.pieslice([1350, 350, 1750, 650], start=120, end=200, fill=(255, 193, 7))

    img.save(SCREENSHOT_PATH, 'PNG')
    print(f'Screenshot created: {SCREENSHOT_PATH}')


def add_text_to_slide(slide, left, top, width, height, text, font_size=18, bold=False, color=None):
    """Helper to add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q3 Sales Performance Review"
    slide1.placeholders[1].text = "Acme Corporation | October 2025"

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    body2 = slide2.placeholders[1]
    tf2 = body2.text_frame
    tf2.text = "Revenue grew 18% YoY to $12.4M in Q3 2025"
    p2a = tf2.add_paragraph()
    p2a.text = "Enterprise segment drove 62% of new bookings"
    p2a.level = 0
    p2b = tf2.add_paragraph()
    p2b.text = "Customer retention rate improved to 94.3%"
    p2b.level = 0
    p2c = tf2.add_paragraph()
    p2c.text = "New product launch contributed $2.1M in pipeline"
    p2c.level = 0
    p2d = tf2.add_paragraph()
    p2d.text = "Operating margin expanded by 3.2 percentage points"
    p2d.level = 0

    # --- Slide 3: Regional Breakdown ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_to_slide(slide3, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                      "Regional Sales Breakdown", font_size=28, bold=True,
                      color=RGBColor(0x1F, 0x38, 0x64))

    table_shape = slide3.shapes.add_table(6, 5, Inches(0.5), Inches(1.3), Inches(12), Inches(4))
    table = table_shape.table
    headers = ["Region", "Q3 Revenue", "Q2 Revenue", "Growth %", "Target Achievement"]
    data = [
        ["North America", "$5,230,000", "$4,410,000", "18.6%", "112%"],
        ["Europe", "$3,180,000", "$2,890,000", "10.0%", "98%"],
        ["Asia Pacific", "$2,450,000", "$1,970,000", "24.4%", "121%"],
        ["Latin America", "$980,000", "$820,000", "19.5%", "105%"],
        ["Middle East & Africa", "$560,000", "$490,000", "14.3%", "93%"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 4: Product Dashboard (TITLE ONLY - NO IMAGE) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_to_slide(slide4, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                      "Product Dashboard", font_size=28, bold=True,
                      color=RGBColor(0x1F, 0x38, 0x64))

    # --- Slide 5: Customer Feedback ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_to_slide(slide5, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                      "Customer Feedback Highlights", font_size=28, bold=True,
                      color=RGBColor(0x1F, 0x38, 0x64))
    quotes = [
        '"The new analytics dashboard has transformed how we track KPIs." - Sarah Chen, VP Operations, TechFlow Inc.',
        '"Response time from the support team dropped from 4 hours to under 30 minutes." - James Rodriguez, CTO, DataBridge Solutions',
        '"We saw a 40% increase in team productivity after adopting the platform." - Priya Patel, Director of Engineering, CloudScale',
    ]
    y_pos = Inches(1.5)
    for quote in quotes:
        add_text_to_slide(slide5, Inches(0.8), y_pos, Inches(11), Inches(1.2),
                          quote, font_size=16, color=RGBColor(0x33, 0x33, 0x33))
        y_pos += Inches(1.6)

    # --- Slide 6: Growth Strategy ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_to_slide(slide6, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                      "Q4 Growth Strategy", font_size=28, bold=True,
                      color=RGBColor(0x1F, 0x38, 0x64))
    strategies = [
        "1. Expand enterprise sales team by 25% in APAC region",
        "2. Launch self-service onboarding portal for SMB segment",
        "3. Introduce tiered pricing with annual commitment discounts",
        "4. Partner integrations with Salesforce, HubSpot, and Slack",
        "5. Invest in AI-powered lead scoring and pipeline forecasting",
    ]
    y_pos = Inches(1.5)
    for s in strategies:
        add_text_to_slide(slide6, Inches(0.8), y_pos, Inches(11), Inches(0.7),
                          s, font_size=18, color=RGBColor(0x33, 0x33, 0x33))
        y_pos += Inches(0.9)

    # --- Slide 7: Next Steps ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_to_slide(slide7, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                      "Next Steps & Action Items", font_size=28, bold=True,
                      color=RGBColor(0x1F, 0x38, 0x64))
    steps = [
        "Finalize Q4 budget allocations by October 15",
        "Schedule regional kick-off meetings for new territory plans",
        "Complete product roadmap review with engineering team",
        "Prepare board presentation with updated financial projections",
        "Launch customer satisfaction survey for enterprise accounts",
    ]
    y_pos = Inches(1.5)
    for step in steps:
        add_text_to_slide(slide7, Inches(0.8), y_pos, Inches(11), Inches(0.7),
                          step, font_size=18, color=RGBColor(0x33, 0x33, 0x33))
        y_pos += Inches(0.9)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Create the screenshot image on the Desktop
    create_screenshot()

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
