"""
Reward Script: Insert product screenshot on slide 4, centered, 8in wide, aspect ratio maintained, 1pt #CCCCCC border
Task ID: impress_sales_028
Domain: libreoffice_impress
Scoring:
  Component 1 — Image present on slide 4            (0.20)
  Component 2 — Image width is 8 inches             (0.20)
  Component 3 — Aspect ratio maintained (16:9)       (0.20)
  Component 4 — Image centered on slide              (0.20)
  Component 5 — 1pt #CCCCCC border around image      (0.20)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_028'


def persist_app_state(domain):
    """Try to save any unsaved LibreOffice state via Ctrl+S."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print(f"PERSIST: ctrl+s sent for {domain}")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # slide 4 (0-indexed)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Find picture shapes on slide 4
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]

    # Component 1: Image present on slide 4 (0.20 points)
    try:
        if len(pictures) > 0:
            print(f"PASS: Component 1 -- Image found on slide 4 ({len(pictures)} picture(s)) (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 -- No images found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    if len(pictures) == 0:
        # No image means all remaining checks fail
        print(f"\nScore: {total_score}/1.0")
        final_score = min(total_score, 1.0)
        print(f"REWARD: {final_score}")
        return final_score

    # Use the first (or only) picture shape
    pic = pictures[0]

    # Component 2: Image width is 8 inches (0.20 points)
    # 8 inches = 7315200 EMU. Allow 2% tolerance.
    try:
        expected_width = Inches(8)  # 7315200 EMU
        actual_width = pic.width
        width_ratio = abs(actual_width - expected_width) / expected_width
        if width_ratio <= 0.02:
            print(f"PASS: Component 2 -- Image width is {actual_width/914400:.3f} inches (expected ~8.0) (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 2 -- Image width is {actual_width/914400:.3f} inches, expected ~8.0 (off by {width_ratio*100:.1f}%)")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Aspect ratio maintained (0.20 points)
    # Original image is 1920x1080 = 16:9 ratio = 1.7778
    # With 8in width, height should be 4.5in. Allow 2% tolerance on ratio.
    try:
        actual_ratio = pic.width / pic.height
        expected_ratio = 1920 / 1080  # 1.7778
        ratio_diff = abs(actual_ratio - expected_ratio) / expected_ratio
        if ratio_diff <= 0.02:
            print(f"PASS: Component 3 -- Aspect ratio {actual_ratio:.4f} matches expected {expected_ratio:.4f} (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 3 -- Aspect ratio {actual_ratio:.4f}, expected {expected_ratio:.4f} (off by {ratio_diff*100:.1f}%)")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Image centered on slide (0.20 points)
    # Center = left + width/2 should match slide_width/2; same for vertical
    # Allow tolerance of 1% of slide dimension
    try:
        pic_center_x = pic.left + pic.width // 2
        pic_center_y = pic.top + pic.height // 2
        slide_center_x = slide_width // 2
        slide_center_y = slide_height // 2

        h_offset = abs(pic_center_x - slide_center_x)
        v_offset = abs(pic_center_y - slide_center_y)
        h_tolerance = slide_width * 0.02
        v_tolerance = slide_height * 0.02

        h_ok = h_offset <= h_tolerance
        v_ok = v_offset <= v_tolerance

        if h_ok and v_ok:
            print(f"PASS: Component 4 -- Image centered (h_off={h_offset} EMU, v_off={v_offset} EMU) (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 4 -- Image not centered (h_off={h_offset/914400:.3f}in, v_off={v_offset/914400:.3f}in)")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: 1pt #CCCCCC border around image (0.20 points)
    # 1pt = 12700 EMU. Allow 10% tolerance on width.
    try:
        line = pic.line
        line_width = line.width
        has_width = False
        has_color = False

        # Check line width: 1pt = 12700 EMU
        if line_width is not None and line_width > 0:
            expected_pt = 12700  # 1pt in EMU
            width_diff = abs(line_width - expected_pt) / expected_pt
            if width_diff <= 0.15:
                has_width = True
                print(f"  Border width: {line_width/12700:.2f} pt (expected 1.0)")
            else:
                print(f"  Border width: {line_width/12700:.2f} pt (expected 1.0, off by {width_diff*100:.0f}%)")
        else:
            print(f"  Border width: None or 0 (expected 1pt)")

        # Check line color: #CCCCCC
        try:
            color_rgb = str(line.color.rgb).upper()
            if color_rgb == 'CCCCCC':
                has_color = True
                print(f"  Border color: #{color_rgb} (correct)")
            else:
                print(f"  Border color: #{color_rgb} (expected #CCCCCC)")
        except Exception:
            print(f"  Border color: not set or not RGB")

        if has_width and has_color:
            print(f"PASS: Component 5 -- 1pt #CCCCCC border present (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 5 -- Border incomplete (width_ok={has_width}, color_ok={has_color})")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
