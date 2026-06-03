"""
Initial Setup: investor_roadshow.pptx — slide 2 with 3 unanimated objects
Task ID: impress_anim_078
Domain: libreoffice_impress

Creates a presentation with:
  Slide 1: Title slide
  Slide 2: Background rectangle + title text box + subtitle text box (NO animations)
  Slide 3: Additional content slide
File saved to ~/Desktop/investor_roadshow.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user/Desktop'
TASK_ID = 'impress_anim_078'
OUTPUT = f'{WORKDIR}/investor_roadshow.pptx'
# Also save with task_id prefix for pipeline contract
OUTPUT_INITIAL = f'{WORKDIR}/{TASK_ID}_initial.pptx'

def create_initial():
    os.makedirs(WORKDIR, exist_ok=True)

    prs = Presentation()
    # Standard widescreen 10x7.5 inches
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title slide ──────────────────────────────────────────────────
    sl1 = prs.slides.add_slide(prs.slide_layouts[0])
    sl1.shapes.title.text = "Investor Roadshow 2025"
    sl1.placeholders[1].text = "Strategic Vision & Growth Opportunities\nQ2 2025"

    # Slide 1 background
    fill1 = sl1.background.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(0x1A, 0x37, 0x6B)

    # ── Slide 2: Investment Thesis — 3 shapes, NONE animated ─────────────────
    sl2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # 2a. Background rectangle covering most of the slide
    bg_rect = sl2.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.5), Inches(0.5), Inches(9.0), Inches(6.5)
    )
    bg_rect.name = "BackgroundRect"
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = RGBColor(0x0A, 0x1E, 0x4A)
    bg_rect.line.fill.background()  # no border

    # 2b. Title text box
    title_box = sl2.shapes.add_textbox(
        Inches(1.0), Inches(1.2), Inches(8.0), Inches(1.4)
    )
    title_box.name = "TitleBox"
    tf_title = title_box.text_frame
    tf_title.word_wrap = False
    p_title = tf_title.paragraphs[0]
    p_title.alignment = PP_ALIGN.LEFT
    run_title = p_title.add_run()
    run_title.text = "Our Investment Thesis"
    run_title.font.name = "Calibri"
    run_title.font.size = Pt(40)
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 2c. Subtitle text box
    sub_box = sl2.shapes.add_textbox(
        Inches(1.0), Inches(2.9), Inches(8.0), Inches(3.6)
    )
    sub_box.name = "SubtitleBox"
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    sub_lines = [
        "• $2.4B total addressable market growing at 18% CAGR",
        "• Proprietary AI platform with 94% customer retention",
        "• Partnerships with 3 Fortune 500 enterprise clients",
        "• Path to profitability in 18 months",
    ]
    for i, line in enumerate(sub_lines):
        if i == 0:
            p = tf_sub.paragraphs[0]
        else:
            p = tf_sub.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.name = "Calibri"
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)

    # ── Slide 3: Financials overview ─────────────────────────────────────────
    sl3 = prs.slides.add_slide(prs.slide_layouts[1])
    sl3.shapes.title.text = "Financial Highlights"
    sl3.placeholders[1].text = (
        "Revenue FY2024: $12.8M (+67% YoY)\n"
        "ARR: $18.5M — Gross Margin: 72%\n"
        "Series B: $35M closed March 2025\n"
        "Runway: 28 months post-raise"
    )

    prs.save(OUTPUT)
    import shutil
    shutil.copy(OUTPUT, OUTPUT_INITIAL)
    print(f"Initial file created: {OUTPUT}")
    print(f"Also saved as: {OUTPUT_INITIAL}")
    print("Slide 2 has 3 shapes: BackgroundRect, TitleBox, SubtitleBox — NO animations")

create_initial()
