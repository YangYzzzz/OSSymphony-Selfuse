"""
FINAL REWARD SCRIPT - SUCCESS
Task: On the very last slide—number 110—I’ve got the word “More” sitting at the bottom. I’d like it to act as a clickable link that jumps all the way back to Slide 1 when viewers press it. How do I set up that hyperlink in LibreOffice Impress?
Generated: 2025-09-10 17:30:13
Status: success
Model: azure-o3
Total Steps: 5
"""

import os
import re
from pptx import Presentation


def verify_hyperlink_to_first_slide(shape, first_slide):
    """Return True if the given shape (or its text runs) contains a hyperlink
    that targets the first slide of the presentation."""

    # 1) Shape-level click action ------------------------------------------------
    try:
        if hasattr(shape, "click_action") and shape.click_action is not None:
            ca = shape.click_action
            if ca.target_slide is not None and ca.target_slide == first_slide:
                # Confirm the action type is a slide jump
                return True
    except Exception as e:
        print("  Error inspecting shape-level click_action:", e)

    # 2) Run-level hyperlink inside the text frame --------------------------------
    try:
        if hasattr(shape, "text_frame") and shape.text_frame is not None:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text and run.text.strip().lower() == "more":
                        h = run.hyperlink

                        # a) Direct address like 'slide1.xml'
                        if h and h.address and re.match(r".*slide1\\.xml$", h.address):
                            return True

                        # b) Relationship target (internal link without address)
                        if h and h._hlinkClick is not None:
                            try:
                                part = run._r.part
                                rel  = part.rels[h._hlinkClick.rId]
                                target_ref = rel.target_ref  # e.g. 'slide1.xml'
                                if re.match(r".*slide1\\.xml$", target_ref):
                                    return True
                            except Exception as e:
                                print("  Error inspecting run hyperlink rel:", e)
    except Exception as e:
        print("  Error inspecting run-level hyperlink:", e)

    return False


def verify_task(file_path):
    """Main verification function.
    Returns a float between 0.0 and 1.0 indicating task completion."""

    print(f"Checking presentation: {file_path}")
    score      = 0.0
    MAX_SCORE  = 1.0

    # -------------------------------------------------------------------------
    # Step 1: Load presentation (no points for loading – prerequisite)
    # -------------------------------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File does not exist")
        return 0.0
    if not file_path.lower().endswith(".pptx"):
        print("✗ File is not a PPTX file")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print("✗ Unable to load PPTX:", e)
        return 0.0

    slide_count = len(prs.slides)
    print(f"Total slides: {slide_count}")

    # -------------------------------------------------------------------------
    # Requirement 1: Presentation has at least 110 slides (0.2 points)
    # -------------------------------------------------------------------------
    if slide_count >= 110:
        print("✓ Slide count requirement met (>=110)")
        score += 0.2
    else:
        print("✗ Insufficient slide count (<110)")

    # Identify first and last slides
    first_slide = prs.slides[0]
    last_slide  = prs.slides[-1]

    # -------------------------------------------------------------------------
    # Requirement 2: "More" text exists on the last slide (0.4 points)
    # -------------------------------------------------------------------------
    more_shape = None
    for shape in last_slide.shapes:
        if hasattr(shape, "text") and shape.text:
            if re.search(r"\bmore\b", shape.text, flags=re.I):
                more_shape = shape
                print(f"✓ Found 'More' text in shape: '{shape.text.strip()}'")
                break

    if more_shape is not None:
        score += 0.4
    else:
        print("✗ Could not find 'More' text on last slide")

    # -------------------------------------------------------------------------
    # Requirement 3: 'More' hyperlink points to first slide (0.4 points)
    # -------------------------------------------------------------------------
    if more_shape is not None:
        if verify_hyperlink_to_first_slide(more_shape, first_slide):
            print("✓ Hyperlink correctly points to first slide")
            score += 0.4
        else:
            print("✗ Hyperlink does not point to first slide")

    # -------------------------------------------------------------------------
    final_score = min(score, MAX_SCORE)
    print(f"Total Score: {final_score}/{MAX_SCORE}")
    return final_score


if __name__ == "__main__":
    FILE = "/home/user/on_the_very_last_slidenumber_110ive_got_the_word_more_sitting_at_the_bottom_id_like_it_to_act_as_a_c_golden.pptx"
    reward = verify_task(FILE)
    print(f"REWARD: {reward}")
