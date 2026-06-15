"""
Initial Setup: Museum Kiosk Presentation - 12 slides, no transitions or timings
Task ID: impress_gf4_026
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_026'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                alignment=PP_ALIGN.LEFT, color=None):
    """Helper to add a textbox with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Museum exhibit themes for 12 slides
    slide_data = [
        {
            "title": "Welcome to the National Science Museum",
            "subtitle": "Interactive Kiosk — Touch to Explore",
            "bg": RGBColor(0x1A, 0x1A, 0x2E),
            "title_color": RGBColor(0xFF, 0xFF, 0xFF),
        },
        {
            "title": "Ancient Civilizations Gallery",
            "body": "Explore artifacts from Mesopotamia, Egypt, and the Indus Valley. "
                    "Our collection spans over 5,000 years of human history, featuring "
                    "cuneiform tablets, pharaonic jewelry, and Harappan seals.",
            "bg": RGBColor(0x2C, 0x1E, 0x0F),
            "title_color": RGBColor(0xD4, 0xA5, 0x74),
        },
        {
            "title": "The Age of Dinosaurs",
            "body": "Walk among life-size reconstructions of Tyrannosaurus Rex, "
                    "Triceratops, and Velociraptor. Our paleontology wing houses "
                    "over 200 fossil specimens dating from the Triassic to Cretaceous periods.",
            "bg": RGBColor(0x0D, 0x2B, 0x0D),
            "title_color": RGBColor(0x7F, 0xC8, 0x7F),
        },
        {
            "title": "Space Exploration Hall",
            "body": "From the Apollo missions to the Mars rovers, discover humanity's "
                    "journey beyond Earth. See authentic mission patches, astronaut suits, "
                    "and a 1:10 scale replica of the International Space Station.",
            "bg": RGBColor(0x0A, 0x0A, 0x2E),
            "title_color": RGBColor(0x87, 0xCE, 0xFA),
        },
        {
            "title": "Ocean Depths Exhibition",
            "body": "Dive into the mysteries of the deep sea. Our immersive exhibit "
                    "showcases bioluminescent creatures, coral reef ecosystems, and the "
                    "technology used to explore the Mariana Trench at 36,000 feet.",
            "bg": RGBColor(0x00, 0x1E, 0x3D),
            "title_color": RGBColor(0x00, 0xBF, 0xFF),
        },
        {
            "title": "Human Body: Inside Out",
            "body": "An interactive journey through anatomy and physiology. Explore "
                    "the cardiovascular system with our 3D heart model, test your reflexes "
                    "at the neuroscience station, and learn about the 206 bones in the "
                    "human skeleton.",
            "bg": RGBColor(0x3D, 0x0C, 0x0C),
            "title_color": RGBColor(0xFF, 0x69, 0x69),
        },
        {
            "title": "Robotics & Artificial Intelligence",
            "body": "Meet our robot guides and watch live demonstrations of machine "
                    "learning algorithms. The AI lab features neural network visualizations "
                    "and hands-on coding stations where visitors program simple robots.",
            "bg": RGBColor(0x1C, 0x1C, 0x2E),
            "title_color": RGBColor(0x00, 0xFF, 0xCC),
        },
        {
            "title": "Climate Science Center",
            "body": "Understanding our changing planet through data and innovation. "
                    "Interactive weather simulations, ice core samples from Antarctica, "
                    "and real-time atmospheric monitoring displays show the impact of "
                    "climate change over 800,000 years.",
            "bg": RGBColor(0x0F, 0x2B, 0x1F),
            "title_color": RGBColor(0x90, 0xEE, 0x90),
        },
        {
            "title": "The Mathematics Playground",
            "body": "Discover the beauty of numbers and patterns. From Fibonacci spirals "
                    "in nature to fractal geometry art, this hands-on gallery makes abstract "
                    "concepts tangible through puzzles and optical illusions.",
            "bg": RGBColor(0x2E, 0x1A, 0x2E),
            "title_color": RGBColor(0xDA, 0x70, 0xD6),
        },
        {
            "title": "Renewable Energy Pavilion",
            "body": "Explore solar, wind, and hydroelectric power through working models. "
                    "Our miniature wind farm generates enough electricity to power the "
                    "exhibit itself, and visitors can race solar-powered cars on our track.",
            "bg": RGBColor(0x1A, 0x2E, 0x1A),
            "title_color": RGBColor(0xFF, 0xD7, 0x00),
        },
        {
            "title": "Photography Through the Ages",
            "body": "From daguerreotypes to digital sensors, trace the evolution of "
                    "capturing light. Our darkroom demonstrations run every hour, and the "
                    "gallery features prize-winning photographs from the Wildlife "
                    "Photographer of the Year competition.",
            "bg": RGBColor(0x1E, 0x1E, 0x1E),
            "title_color": RGBColor(0xFF, 0xFF, 0xFF),
        },
        {
            "title": "Thank You for Visiting",
            "subtitle": "We hope you enjoyed your journey through science and discovery.\n"
                        "Visit our gift shop on Level 1.\n"
                        "Open daily 9:00 AM — 6:00 PM",
            "bg": RGBColor(0x1A, 0x1A, 0x2E),
            "title_color": RGBColor(0xFF, 0xFF, 0xFF),
        },
    ]

    for i, data in enumerate(slide_data):
        if i == 0 or i == 11:
            # Title slides (first and last)
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only

        # Set background
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = data["bg"]

        # Title
        add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.5),
                    data["title"], font_size=36, bold=True,
                    alignment=PP_ALIGN.CENTER, color=data["title_color"])

        # Body or subtitle
        if "subtitle" in data:
            add_textbox(slide, Inches(2), Inches(3), Inches(9), Inches(2.5),
                        data["subtitle"], font_size=24, bold=False,
                        alignment=PP_ALIGN.CENTER, color=RGBColor(0xCC, 0xCC, 0xCC))
        elif "body" in data:
            add_textbox(slide, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4),
                        data["body"], font_size=20, bold=False,
                        alignment=PP_ALIGN.LEFT, color=RGBColor(0xDD, 0xDD, 0xDD))

        # Add slide number indicator
        add_textbox(slide, Inches(6), Inches(6.8), Inches(1.3), Inches(0.5),
                    f"{i+1} / 12", font_size=12, bold=False,
                    alignment=PP_ALIGN.CENTER, color=RGBColor(0x88, 0x88, 0x88))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
