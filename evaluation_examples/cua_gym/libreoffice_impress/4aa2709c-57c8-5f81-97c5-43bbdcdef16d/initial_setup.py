"""
Initial Setup: Sales pitch presentation with 8 slides, no transitions
Task ID: impress_sales_026
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_026'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=None, font_name="Arial"):
    """Helper to add a text box with common formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tf


def add_bullet_list(text_frame, items, font_size=14, color=None, font_name="Arial"):
    """Add bullet items to an existing text frame."""
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            if color:
                run.font.color.rgb = color


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color palette
    DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
    MEDIUM_BLUE = RGBColor(0x2E, 0x86, 0xC1)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0x95, 0xA5, 0xA6)
    DARK_GRAY = RGBColor(0x2C, 0x3E, 0x50)
    ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    add_textbox(slide1, Inches(1.5), Inches(1.8), Inches(10), Inches(1.5),
                "Smooth Pitch", font_size=44, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(1.5), Inches(3.5), Inches(10), Inches(1),
                "Revolutionizing Sales Enablement for Enterprise Teams",
                font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(1.5), Inches(5.5), Inches(10), Inches(0.6),
                "Q2 2025 Investor Presentation  |  Confidential",
                font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ---- Slide 2: Agenda / Overview ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide2, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Agenda", font_size=36, bold=True, color=DARK_BLUE,
                alignment=PP_ALIGN.LEFT)
    tf2 = add_textbox(slide2, Inches(1.2), Inches(1.8), Inches(10), Inches(4.5),
                      "1. Market Landscape & Opportunity", font_size=20, color=DARK_GRAY)
    items2 = [
        "2. The Problem We Solve",
        "3. Our Solution: Smooth Pitch Platform",
        "4. Product Demo & Key Features",
        "5. Market Traction & Revenue Growth",
        "6. Competitive Positioning",
        "7. Pricing & Business Model",
        "8. Next Steps & Contact Information",
    ]
    for item in items2:
        p = tf2.add_paragraph()
        p.text = item
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(20)
            run.font.color.rgb = DARK_GRAY

    # ---- Slide 3: Problem Statement ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide3, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "The Problem", font_size=36, bold=True, color=DARK_BLUE)

    problems = [
        ("68%", "of sales reps miss their quarterly targets due to inconsistent messaging"),
        ("$2.3M", "average annual revenue lost per enterprise from poor pitch materials"),
        ("14 hrs/week", "spent by sales teams manually updating presentations"),
        ("42%", "of prospects disengage when pitches lack personalization"),
    ]
    for i, (stat, desc) in enumerate(problems):
        y_pos = Inches(1.8 + i * 1.3)
        add_textbox(slide3, Inches(1.0), y_pos, Inches(2), Inches(0.8),
                    stat, font_size=28, bold=True, color=MEDIUM_BLUE)
        add_textbox(slide3, Inches(3.2), y_pos + Inches(0.1), Inches(8.5), Inches(0.8),
                    desc, font_size=18, color=DARK_GRAY)

    # ---- Slide 4: Our Solution ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    fill4 = slide4.background.fill
    fill4.solid()
    fill4.fore_color.rgb = DARK_BLUE

    add_textbox(slide4, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Our Solution: Smooth Pitch", font_size=36, bold=True, color=WHITE)
    add_textbox(slide4, Inches(0.8), Inches(1.8), Inches(11), Inches(1),
                "AI-powered sales enablement that creates personalized, data-driven "
                "presentations in minutes, not hours.",
                font_size=20, color=LIGHT_GRAY)

    features_sol = [
        "Smart Template Engine - Auto-generates tailored pitch decks from CRM data",
        "Real-Time Analytics - Track prospect engagement and optimize messaging",
        "Brand Consistency - Centralized asset library with version control",
        "Seamless Integration - Works with Salesforce, HubSpot, and 40+ tools",
    ]
    for i, feat in enumerate(features_sol):
        add_textbox(slide4, Inches(1.2), Inches(3.2 + i * 0.9), Inches(10.5), Inches(0.7),
                    feat, font_size=16, color=WHITE)

    # ---- Slide 5: Product Features ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide5, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Key Features", font_size=36, bold=True, color=DARK_BLUE)

    # Left column
    left_features = [
        ("AI Content Generation", "Generate compelling copy tailored to each prospect's industry and pain points"),
        ("Dynamic Data Visualization", "Auto-create charts and infographics from your sales data"),
        ("Collaboration Hub", "Real-time co-editing with role-based permissions and approval workflows"),
    ]
    for i, (title, desc) in enumerate(left_features):
        y = Inches(1.8 + i * 1.6)
        add_textbox(slide5, Inches(0.8), y, Inches(5.5), Inches(0.5),
                    title, font_size=18, bold=True, color=MEDIUM_BLUE)
        add_textbox(slide5, Inches(0.8), y + Inches(0.5), Inches(5.5), Inches(0.8),
                    desc, font_size=14, color=DARK_GRAY)

    # Right column
    right_features = [
        ("Performance Analytics", "A/B test slides, track viewing time, and measure conversion rates"),
        ("Mobile-First Design", "Present from any device with offline capability and gesture controls"),
        ("Enterprise Security", "SOC 2 Type II certified, SSO, and advanced access controls"),
    ]
    for i, (title, desc) in enumerate(right_features):
        y = Inches(1.8 + i * 1.6)
        add_textbox(slide5, Inches(7.0), y, Inches(5.5), Inches(0.5),
                    title, font_size=18, bold=True, color=MEDIUM_BLUE)
        add_textbox(slide5, Inches(7.0), y + Inches(0.5), Inches(5.5), Inches(0.8),
                    desc, font_size=14, color=DARK_GRAY)

    # ---- Slide 6: Market Analysis ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide6, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Market Traction", font_size=36, bold=True, color=DARK_BLUE)

    metrics = [
        ("$4.8M ARR", "Annual Recurring Revenue\n+187% YoY Growth"),
        ("2,400+", "Enterprise Users\nacross 18 countries"),
        ("98.2%", "Customer Retention Rate\nIndustry avg: 85%"),
        ("4.8/5.0", "G2 Rating\n320+ verified reviews"),
    ]
    for i, (metric, detail) in enumerate(metrics):
        x = Inches(0.8 + i * 3.1)
        add_textbox(slide6, x, Inches(2.0), Inches(2.8), Inches(0.8),
                    metric, font_size=28, bold=True, color=ACCENT_GREEN,
                    alignment=PP_ALIGN.CENTER)
        tf_m = add_textbox(slide6, x, Inches(3.0), Inches(2.8), Inches(1.5),
                           detail.split('\n')[0], font_size=14, color=DARK_GRAY,
                           alignment=PP_ALIGN.CENTER)
        if '\n' in detail:
            p2 = tf_m.add_paragraph()
            p2.text = detail.split('\n')[1]
            p2.alignment = PP_ALIGN.CENTER
            for run in p2.runs:
                run.font.size = Pt(12)
                run.font.color.rgb = LIGHT_GRAY

    # Notable clients
    add_textbox(slide6, Inches(0.8), Inches(5.0), Inches(11), Inches(0.6),
                "Trusted by: Deloitte  |  Accenture  |  Siemens  |  Philips  |  Stripe  |  Datadog",
                font_size=16, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # ---- Slide 7: Pricing ----
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide7, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Pricing & Plans", font_size=36, bold=True, color=DARK_BLUE)

    plans = [
        ("Starter", "$29/user/mo", ["Up to 10 users", "50 presentations/month",
                                     "Basic analytics", "Email support"]),
        ("Professional", "$79/user/mo", ["Up to 100 users", "Unlimited presentations",
                                          "Advanced analytics & A/B testing",
                                          "Priority support", "CRM integrations"]),
        ("Enterprise", "Custom", ["Unlimited users", "Custom AI training",
                                   "Dedicated success manager", "SSO & SCIM",
                                   "Custom SLA & on-premise option"]),
    ]
    for i, (name, price, features) in enumerate(plans):
        x = Inches(0.8 + i * 4.0)
        add_textbox(slide7, x, Inches(1.8), Inches(3.5), Inches(0.6),
                    name, font_size=24, bold=True, color=MEDIUM_BLUE,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide7, x, Inches(2.5), Inches(3.5), Inches(0.6),
                    price, font_size=20, bold=True, color=DARK_GRAY,
                    alignment=PP_ALIGN.CENTER)
        for j, feat in enumerate(features):
            add_textbox(slide7, x + Inches(0.3), Inches(3.3 + j * 0.6),
                        Inches(3.0), Inches(0.5),
                        feat, font_size=13, color=DARK_GRAY)

    # ---- Slide 8: Contact / CTA ----
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    fill8 = slide8.background.fill
    fill8.solid()
    fill8.fore_color.rgb = DARK_BLUE

    add_textbox(slide8, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
                "Ready to Transform Your Sales?",
                font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide8, Inches(1.5), Inches(3.2), Inches(10), Inches(1),
                "Schedule a personalized demo and see how Smooth Pitch can "
                "accelerate your revenue growth.",
                font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    contact_lines = [
        "Elena Vasquez, VP of Sales",
        "elena.vasquez@smoothpitch.io  |  +1 (415) 555-0198",
        "www.smoothpitch.io/demo",
    ]
    for i, line in enumerate(contact_lines):
        add_textbox(slide8, Inches(1.5), Inches(4.8 + i * 0.6), Inches(10), Inches(0.5),
                    line, font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
