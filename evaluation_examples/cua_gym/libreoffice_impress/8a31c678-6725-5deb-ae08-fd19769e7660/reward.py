"""
Reward Script: Set slide backgrounds and speaker notes for slides 2 and 4
Task ID: osworld_impress_note_bg_combined_007
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 2 has light blue background (#ADD8E6)       - 0.25 pts
  Component 2: Slide 2 has correct speaker note                   - 0.25 pts
  Component 3: Slide 4 has light green background (#90EE90)       - 0.25 pts
  Component 4: Slide 4 has correct speaker note                   - 0.25 pts
  Total: 1.0
"""

import os

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_note_bg_combined_007'

# Expected values derived from task_config.json context
SLIDE2_EXPECTED_BG = 'ADD8E6'   # light blue
SLIDE2_EXPECTED_NOTE = 'Speaker 1: Introduce the problem statement (3 minutes)'

SLIDE4_EXPECTED_BG = '90EE90'   # light green
SLIDE4_EXPECTED_NOTE = 'Speaker 2: Present the solution demo (5 minutes)'


def get_slide_background_hex(slide):
    """Return the solid background color as a hex string (e.g. 'ADD8E6'), or None."""
    try:
        fill = slide.background.fill
        if fill.type == 1:  # MSO_FILL.SOLID
            return str(fill.fore_color.rgb).upper()
        elif fill.type == 5:  # inherited from master
            master_fill = slide.slide_layout.slide_master.background.fill
            if master_fill.type == 1:
                return str(master_fill.fore_color.rgb).upper()
    except Exception:
        pass
    return None


def get_slide_notes(slide):
    """Return the notes text stripped of whitespace, or empty string."""
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"CRITICAL: Expected at least 4 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]  # 0-indexed: slide 2 = index 1
    slide4 = prs.slides[3]  # 0-indexed: slide 4 = index 3

    # Component 1: Slide 2 has light blue background (#ADD8E6) (0.25 points)
    try:
        actual_bg2 = get_slide_background_hex(slide2)
        if actual_bg2 == SLIDE2_EXPECTED_BG:
            print(f"PASS: Component 1 — Slide 2 background is light blue #{SLIDE2_EXPECTED_BG} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — Slide 2 background expected #{SLIDE2_EXPECTED_BG}, found {actual_bg2}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 2 has the correct speaker note (0.25 points)
    try:
        actual_note2 = get_slide_notes(slide2)
        if actual_note2 == SLIDE2_EXPECTED_NOTE:
            print(f"PASS: Component 2 — Slide 2 note matches expected (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 — Slide 2 note expected {repr(SLIDE2_EXPECTED_NOTE)}, found {repr(actual_note2)}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 4 has light green background (#90EE90) (0.25 points)
    try:
        actual_bg4 = get_slide_background_hex(slide4)
        if actual_bg4 == SLIDE4_EXPECTED_BG:
            print(f"PASS: Component 3 — Slide 4 background is light green #{SLIDE4_EXPECTED_BG} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 — Slide 4 background expected #{SLIDE4_EXPECTED_BG}, found {actual_bg4}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 4 has the correct speaker note (0.25 points)
    try:
        actual_note4 = get_slide_notes(slide4)
        if actual_note4 == SLIDE4_EXPECTED_NOTE:
            print(f"PASS: Component 4 — Slide 4 note matches expected (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 4 — Slide 4 note expected {repr(SLIDE4_EXPECTED_NOTE)}, found {repr(actual_note4)}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the VM environment
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
