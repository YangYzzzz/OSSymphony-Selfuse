"""
Initial Setup: Panel discussion deck - 6 slides, white backgrounds, no notes on slides 2 and 4
Task ID: osworld_impress_note_bg_combined_007
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_note_bg_combined_007'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide data for a 6-slide panel discussion deck
    slides_data = [
        {
            "layout_idx": 0,  # Title Slide
            "title": "Global Climate Policy: A Panel Discussion",
            "subtitle": "Annual Summit on Sustainable Development | March 2025",
        },
        {
            "layout_idx": 1,  # Title + Content
            "title": "The Problem Statement",
            "content": (
                "• Global CO₂ emissions reached 36.8 Gt in 2024\n"
                "• Temperature anomalies exceed 1.2°C above pre-industrial levels\n"
                "• 40% of countries lack enforceable climate legislation\n"
                "• Funding gap for adaptation: $300 billion annually"
            ),
        },
        {
            "layout_idx": 1,
            "title": "Current Policy Landscape",
            "content": (
                "• Paris Agreement: 196 signatories with varied commitments\n"
                "• EU Green Deal targets 55% reduction by 2030\n"
                "• US IRA: $369B in clean energy investment through 2032\n"
                "• Developing nations still rely on fossil fuels for 78% of energy"
            ),
        },
        {
            "layout_idx": 1,
            "title": "Proposed Solutions & Innovations",
            "content": (
                "• Carbon pricing mechanisms: cap-and-trade vs. carbon tax\n"
                "• Green hydrogen infrastructure rollout by 2035\n"
                "• Direct Air Capture technology scaling roadmap\n"
                "• Public-private partnership model for renewable financing"
            ),
        },
        {
            "layout_idx": 1,
            "title": "Panel Discussion: Key Questions",
            "content": (
                "1. How can nations align on binding emission targets?\n"
                "2. What role should the private sector play in funding?\n"
                "3. How do we balance economic growth with sustainability?\n"
                "4. What are the most scalable technologies for emerging economies?"
            ),
        },
        {
            "layout_idx": 0,
            "title": "Q&A and Closing Remarks",
            "subtitle": "Thank you — Open Floor for Questions",
        },
    ]

    for i, sd in enumerate(slides_data):
        layout = prs.slide_layouts[sd["layout_idx"]]
        slide = prs.slides.add_slide(layout)

        # Set white background explicitly
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = sd["title"]

        # Set content/subtitle
        if sd["layout_idx"] == 0:
            # Title Slide — use placeholder index 1 for subtitle
            try:
                slide.placeholders[1].text = sd.get("subtitle", "")
            except (KeyError, IndexError):
                pass
        else:
            # Title + Content — use placeholder index 1 for content body
            try:
                slide.placeholders[1].text = sd.get("content", "")
            except (KeyError, IndexError):
                pass

        # Slides 2 and 4 (index 1 and 3) must have NO notes
        # (do not access notes_slide property to avoid inadvertently creating them)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
