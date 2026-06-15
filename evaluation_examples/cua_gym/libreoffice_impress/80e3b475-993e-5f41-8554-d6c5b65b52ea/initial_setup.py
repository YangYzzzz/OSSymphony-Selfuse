"""
Initial Setup: University presentation template with 6 slides, no custom branding on slide master.
Task ID: impress_teach_028
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_028'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "State University"
    slide1.placeholders[1].text = "Annual Academic Report 2025-2026"

    # --- Slide 2: Mission & Vision ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Mission & Vision"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Empowering students through transformative education and groundbreaking research"
    p2 = body2.add_paragraph()
    p2.text = "Our commitment to academic excellence drives innovation across 12 colleges and 180+ programs"
    p2.level = 0
    p3 = body2.add_paragraph()
    p3.text = "Serving over 28,000 students from 50 states and 90 countries"
    p3.level = 0

    # --- Slide 3: Enrollment Statistics ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Enrollment Statistics"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Fall 2025 Enrollment Highlights"
    entries = [
        "Undergraduate: 21,450 students (+3.2% from 2024)",
        "Graduate: 5,890 students (+1.8% from 2024)",
        "International: 4,200 students representing 90 countries",
        "First-Year Retention Rate: 91.3%",
        "Four-Year Graduation Rate: 72.8%",
    ]
    for entry in entries:
        p = body3.add_paragraph()
        p.text = entry
        p.level = 1

    # --- Slide 4: Research Highlights ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Research Highlights"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Fiscal Year 2025 Research Expenditures: $342 Million"
    highlights = [
        "National Science Foundation grants: $48.5M across 127 projects",
        "New biomedical research facility opening Fall 2026",
        "12 faculty members elected to National Academies",
        "Patent filings increased 22% year-over-year",
    ]
    for h in highlights:
        p = body4.add_paragraph()
        p.text = h
        p.level = 1

    # --- Slide 5: Campus Development ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Campus Development"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Major Capital Projects 2025-2028"
    projects = [
        "STEM Innovation Center: $95M, completion Spring 2027",
        "Student Wellness Complex: $42M, breaking ground Fall 2025",
        "Sustainability Initiative: Carbon neutral campus by 2030",
    ]
    for proj in projects:
        p = body5.add_paragraph()
        p.text = proj
        p.level = 1

    # --- Slide 6: Thank You ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[0])
    slide6.shapes.title.text = "Thank You"
    slide6.placeholders[1].text = "Questions? Contact: provost@stateu.edu"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
