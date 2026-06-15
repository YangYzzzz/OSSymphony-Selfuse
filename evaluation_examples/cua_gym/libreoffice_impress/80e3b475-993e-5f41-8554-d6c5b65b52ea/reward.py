"""
Reward Script: Add university branding rectangle to slide master
Task ID: impress_teach_028
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): Rectangle shape exists on slide master with correct size (~1x0.5 in)
  Component 2 (0.3): Rectangle at top-left with fill color #1A237E
  Component 3 (0.4): Text 'STATE U' in white, 10pt, bold
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_028'


def is_approximately_equal(val1, val2, tolerance=0.05):
    """Check if two values are approximately equal within tolerance."""
    if val1 == val2:
        return True
    if val1 == 0 or val2 == 0:
        return abs(val1 - val2) <= tolerance * max(abs(val1), abs(val2), 1)
    return abs(val1 - val2) / max(abs(val1), abs(val2)) <= tolerance


def find_brand_rect(master):
    """Find a non-placeholder rectangle shape on the slide master that contains 'STATE U'."""
    for shape in master.shapes:
        # Skip placeholders - we want only added shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            continue
        # Check if it's an auto shape (rectangle) with text
        if hasattr(shape, 'text') and 'STATE' in shape.text.upper():
            return shape
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Get the first slide master
    if len(prs.slide_masters) == 0:
        print("FAIL: No slide masters found")
        print("REWARD: 0.0")
        return 0.0

    master = prs.slide_masters[0]
    brand_shape = find_brand_rect(master)

    if brand_shape is None:
        print("FAIL: No branding rectangle found on slide master")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Found candidate shape: name={brand_shape.name}, type={brand_shape.shape_type}")
    print(f"INFO: Position=({brand_shape.left},{brand_shape.top}), Size=({brand_shape.width},{brand_shape.height})")
    print(f"INFO: Text={repr(brand_shape.text)}")

    # Component 1: Rectangle shape exists on slide master with correct size (~1x0.5 inches) (0.3 points)
    try:
        expected_width = Inches(1)   # 914400 EMU
        expected_height = Inches(0.5)  # 457200 EMU
        width_ok = is_approximately_equal(brand_shape.width, expected_width, tolerance=0.05)
        height_ok = is_approximately_equal(brand_shape.height, expected_height, tolerance=0.05)

        if width_ok and height_ok:
            print(f"PASS: Component 1 -- Rectangle size correct: {brand_shape.width}x{brand_shape.height} EMU (~1x0.5 in) (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 -- Expected ~914400x457200 EMU, found {brand_shape.width}x{brand_shape.height}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Rectangle at top-left with fill color #1A237E (0.3 points)
    try:
        # Check position: top-left corner means left~0 and top~0
        pos_ok = (brand_shape.left <= Inches(0.2) and brand_shape.top <= Inches(0.2))

        # Check fill color
        fill_ok = False
        try:
            fill = brand_shape.fill
            if fill.type is not None and fill.type == 1:  # SOLID
                color_str = str(fill.fore_color.rgb).upper()
                fill_ok = (color_str == '1A237E')
                print(f"INFO: Fill color = {color_str}")
            else:
                print(f"INFO: Fill type = {fill.type} (expected SOLID=1)")
        except Exception as e:
            print(f"INFO: Fill check error: {e}")

        if pos_ok and fill_ok:
            print(f"PASS: Component 2 -- Top-left position and fill #1A237E correct (0.3 pts)")
            total_score += 0.3
        elif pos_ok:
            print(f"FAIL: Component 2 -- Position OK but fill color wrong")
            total_score += 0.1
        elif fill_ok:
            print(f"FAIL: Component 2 -- Fill color OK but position wrong (left={brand_shape.left}, top={brand_shape.top})")
            total_score += 0.1
        else:
            print(f"FAIL: Component 2 -- Both position and fill wrong")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Text 'STATE U' in white, 10pt, bold (0.4 points)
    try:
        text_ok = False
        font_ok = False
        color_ok = False

        # Check text content
        shape_text = brand_shape.text.strip()
        text_ok = (shape_text == 'STATE U')
        if text_ok:
            print(f"INFO: Text matches 'STATE U'")
        else:
            print(f"INFO: Text is {repr(shape_text)}, expected 'STATE U'")

        # Check font properties on runs
        if hasattr(brand_shape, 'text_frame'):
            for para in brand_shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        # Check bold
                        is_bold = (run.font.bold is True)
                        # Check size (~10pt = 127000 EMU)
                        is_10pt = (run.font.size is not None and
                                   is_approximately_equal(run.font.size, Pt(10), tolerance=0.05))
                        # Check white color
                        is_white = False
                        try:
                            if run.font.color.type is not None:
                                rgb_str = str(run.font.color.rgb).upper()
                                is_white = (rgb_str == 'FFFFFF')
                                print(f"INFO: Font color = {rgb_str}")
                        except:
                            pass

                        print(f"INFO: bold={run.font.bold}, size={run.font.size}, white={is_white}")

                        font_ok = (is_bold and is_10pt)
                        color_ok = is_white

        sub_score = 0.0
        if text_ok:
            sub_score += 0.15
        if font_ok:
            sub_score += 0.15
        if color_ok:
            sub_score += 0.1

        if text_ok and font_ok and color_ok:
            print(f"PASS: Component 3 -- Text 'STATE U', white, 10pt bold all correct (0.4 pts)")
        else:
            details = f"text={'OK' if text_ok else 'FAIL'}, bold+10pt={'OK' if font_ok else 'FAIL'}, white={'OK' if color_ok else 'FAIL'}"
            print(f"PARTIAL: Component 3 -- {details} ({sub_score} pts)")

        total_score += sub_score
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
