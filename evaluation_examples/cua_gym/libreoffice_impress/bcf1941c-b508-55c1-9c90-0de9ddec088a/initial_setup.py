#!/usr/bin/env python3
"""
initial_setup.py for impress_tct_091
Creates Warnings_Slide.pptx with 4 slides. Slide 3 has IMPORTANT (x2) and CRITICAL (x1)
in plain black text, NO highlights.
"""
import os
import subprocess
import shlex
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

OUTPUT_PATH = "/home/user/Warnings_Slide.pptx"

prs = Presentation()

# --- Slide 1: Title slide ---
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
slide1.shapes.title.text = "Safety Procedures Manual"
slide1.placeholders[1].text = "Workplace Safety & Compliance Department"

# --- Slide 2: General procedures ---
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
slide2.shapes.title.text = "General Safety Procedures"
body2 = slide2.placeholders[1].text_frame
body2.text = "All employees must complete safety orientation before beginning work."
p2 = body2.add_paragraph()
p2.text = "Regular safety audits are conducted on a quarterly basis."
p2.level = 0
p3 = body2.add_paragraph()
p3.text = "Report any hazards or unsafe conditions to your team lead immediately."
p3.level = 0
p4 = body2.add_paragraph()
p4.text = "Emergency exits must remain clear and accessible at all times."
p4.level = 0

# --- Slide 3: Key slide with IMPORTANT and CRITICAL ---
slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
txBox = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.0), Inches(6.0))
tf = txBox.text_frame
tf.word_wrap = True

# Helper to add a run with black font
def add_run(paragraph, text, bold=False, size=Pt(16)):
    run = paragraph.add_run()
    run.text = text
    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    run.font.size = size
    run.font.bold = bold
    run.font.name = "Arial"
    return run

# Title for slide 3
p_title = tf.paragraphs[0]
p_title.alignment = PP_ALIGN.LEFT
add_run(p_title, "Maintenance Safety Guidelines", bold=True, size=Pt(24))

# Paragraph 1: contains first IMPORTANT
p1 = tf.add_paragraph()
p1.alignment = PP_ALIGN.LEFT
p1.space_before = Pt(12)
add_run(p1, "Before proceeding with any maintenance work, it is ")
add_run(p1, "IMPORTANT", bold=True)
add_run(p1, " to review all safety protocols. Failure to follow proper procedures could result in equipment damage.")

# Paragraph 2: contains CRITICAL
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.LEFT
p2.space_before = Pt(12)
add_run(p2, "Note that ")
add_run(p2, "CRITICAL", bold=True)
add_run(p2, " safety violations must be reported immediately to your supervisor. Delayed reporting can lead to severe consequences.")

# Paragraph 3: contains second IMPORTANT
p3 = tf.add_paragraph()
p3.alignment = PP_ALIGN.LEFT
p3.space_before = Pt(12)
add_run(p3, "Additionally, it is ")
add_run(p3, "IMPORTANT", bold=True)
add_run(p3, " to wear appropriate personal protective equipment at all times during operations.")

# Paragraph 4: plain paragraph
p4 = tf.add_paragraph()
p4.alignment = PP_ALIGN.LEFT
p4.space_before = Pt(12)
add_run(p4, "Please consult the full safety handbook for detailed procedures and checklists.")

# --- Slide 4: Summary ---
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
slide4.shapes.title.text = "Summary & Next Steps"
body4 = slide4.placeholders[1].text_frame
body4.text = "Review all highlighted safety guidelines before starting work."
bp1 = body4.add_paragraph()
bp1.text = "Attend mandatory safety training sessions."
bp1.level = 0
bp2 = body4.add_paragraph()
bp2.text = "Keep documentation updated and accessible."
bp2.level = 0
bp3 = body4.add_paragraph()
bp3.text = "Contact the Safety Office for any questions or concerns."
bp3.level = 0

prs.save(OUTPUT_PATH)
print(f"Saved presentation to {OUTPUT_PATH}")

# Launch LibreOffice Impress
env = os.environ.copy()
env["DISPLAY"] = ":0"
subprocess.Popen(
    shlex.split(f'libreoffice --impress "{OUTPUT_PATH}"'),
    stdout=subprocess.DEVNULL,
    stderr=subprocess.DEVNULL,
    env=env,
)
time.sleep(2)
print("LibreOffice Impress launched.")
