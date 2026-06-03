#!/usr/bin/env python3
"""
Reward script for impress_tct_091:
Check that 'IMPORTANT' (x2) and 'CRITICAL' (x1) on slide 3 have yellow (#FFFF00) highlights,
and no other text on slide 3 has highlighting.
"""

import sys
from pptx import Presentation
from pptx.oxml.ns import qn

PPTX_PATH = "/home/user/Warnings_Slide.pptx"
TARGET_WORDS = {"IMPORTANT", "CRITICAL"}
EXPECTED_COLOR = "FFFF00"


def get_highlight_color(run):
    """Return the highlight srgbClr val for a run, or None if no highlight."""
    rPr = run._r.find(qn('a:rPr'))
    if rPr is None:
        return None
    hl = rPr.find(qn('a:highlight'))
    if hl is None:
        return None
    clr = hl.find(qn('a:srgbClr'))
    if clr is not None:
        return clr.get('val')
    return None


def main():
    score = 0.0

    try:
        prs = Presentation(PPTX_PATH)
    except Exception as e:
        print(f"Error opening presentation: {e}")
        print(f"REWARD: 0.0")
        return

    if len(prs.slides) < 3:
        print("Presentation has fewer than 3 slides.")
        print(f"REWARD: 0.0")
        return

    slide = prs.slides[2]  # slide 3 (0-indexed)

    # Collect all runs from all text shapes on slide 3
    target_runs_found = []  # (text, highlighted_correctly)
    other_runs_highlighted = 0
    total_other_runs = 0

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                text = (run.text or "").strip()
                if not text:
                    continue
                hl_color = get_highlight_color(run)
                if text in TARGET_WORDS:
                    is_correct = (hl_color is not None and hl_color.upper() == EXPECTED_COLOR)
                    target_runs_found.append((text, is_correct))
                else:
                    total_other_runs += 1
                    if hl_color is not None:
                        other_runs_highlighted += 1

    # Component 1 (0.3): First IMPORTANT has yellow highlight
    important_runs = [ok for (t, ok) in target_runs_found if t == "IMPORTANT"]
    if len(important_runs) >= 1 and important_runs[0]:
        score += 0.3

    # Component 2 (0.3): Second IMPORTANT has yellow highlight
    if len(important_runs) >= 2 and important_runs[1]:
        score += 0.3

    # Component 3 (0.2): CRITICAL has yellow highlight
    critical_runs = [ok for (t, ok) in target_runs_found if t == "CRITICAL"]
    if len(critical_runs) >= 1 and critical_runs[0]:
        score += 0.2

    # Component 4 (0.2): No other text on slide 3 has highlighting
    # Only award if at least one target word is correctly highlighted (avoid free points on unmodified file)
    targets_highlighted = sum(1 for (_, ok) in target_runs_found if ok)
    if other_runs_highlighted == 0 and targets_highlighted > 0:
        score += 0.2

    score = round(min(score, 1.0), 1)
    print(f"Target runs found: {len(target_runs_found)} (IMPORTANT: {len(important_runs)}, CRITICAL: {len(critical_runs)})")
    print(f"Other runs with highlighting: {other_runs_highlighted}/{total_other_runs}")
    print(f"REWARD: {score}")


if __name__ == "__main__":
    main()
