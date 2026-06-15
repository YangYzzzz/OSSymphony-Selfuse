"""
Reward Script: Remove broken OLE object from slide 6 and insert data_table.png image
Task ID: impress_fix_053
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Broken OLE shapes removed from slide 6
  Component 2 (0.4): Picture inserted on slide 6 matching data_table.png
  Component 3 (0.2): Inserted image has reasonable size/position
"""

import os
import hashlib
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_053'
DATA_TABLE_PATH = '/home/user/Desktop/data_table.png'


def persist_app_state(domain):
    """Attempt to save any unsaved LibreOffice edits via Ctrl+S."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"CRITICAL: Expected at least 6 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide6 = prs.slides[5]  # 0-indexed, slide 6

    # Load reference image hash for comparison
    try:
        with open(DATA_TABLE_PATH, 'rb') as f:
            ref_blob = f.read()
        ref_md5 = hashlib.md5(ref_blob).hexdigest()
        print(f"INFO: Reference data_table.png md5={ref_md5}, size={len(ref_blob)}")
    except Exception as e:
        print(f"WARN: Could not load reference image {DATA_TABLE_PATH}: {e}")
        ref_md5 = None
        ref_blob = None

    # Collect shape info from slide 6
    shape_names = []
    broken_ole_names = {'Broken_OLE_Object', 'Broken_OLE_XMark', 'Broken_OLE_Label'}
    found_broken = []
    picture_shapes = []

    for shape in slide6.shapes:
        shape_names.append(shape.name)
        if shape.name in broken_ole_names:
            found_broken.append(shape.name)
        # Also check for shapes with "Cannot be displayed" text (in case renamed)
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if 'Cannot be displayed' in text or (text == 'X' and shape.width > Inches(1)):
                if shape.name not in broken_ole_names:
                    found_broken.append(f"{shape.name} (text indicates broken OLE)")
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture_shapes.append(shape)

    print(f"INFO: Slide 6 shapes: {shape_names}")

    # Component 1: Broken OLE shapes removed (0.4 points)
    # In initial_env, slide 6 has Broken_OLE_Object, Broken_OLE_XMark, Broken_OLE_Label
    # These must be gone in the golden/agent result
    try:
        if len(found_broken) == 0:
            print(f"PASS: Component 1 — All broken OLE shapes removed from slide 6 (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 — Broken OLE shapes still present: {found_broken}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Picture shape with data_table.png content exists on slide 6 (0.4 points)
    # In initial_env, there are NO picture shapes on slide 6
    try:
        if len(picture_shapes) == 0:
            print(f"FAIL: Component 2 — No picture shapes found on slide 6")
        else:
            image_matched = False
            for pic in picture_shapes:
                try:
                    pic_blob = pic.image.blob
                    pic_md5 = hashlib.md5(pic_blob).hexdigest()
                    print(f"INFO: Found picture '{pic.name}' md5={pic_md5}, size={len(pic_blob)}")

                    if ref_md5 is not None and pic_md5 == ref_md5:
                        image_matched = True
                        break
                    elif ref_blob is not None and len(pic_blob) == len(ref_blob):
                        # Size match as fallback (image may have been re-encoded)
                        image_matched = True
                        break
                except Exception as e:
                    print(f"WARN: Could not read image blob from {pic.name}: {e}")

            if image_matched:
                print(f"PASS: Component 2 — data_table.png image found on slide 6 (0.4 pts)")
                total_score += 0.4
            else:
                # Partial credit if there's any picture (might be a screenshot of data_table)
                print(f"PARTIAL: Component 2 — Picture found but doesn't match data_table.png exactly")
                total_score += 0.2
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Inserted image has reasonable size and position (0.2 points)
    # The image should be reasonably large (not tiny) and within slide bounds
    try:
        if len(picture_shapes) > 0:
            pic = picture_shapes[0]
            slide_width = prs.slide_width
            slide_height = prs.slide_height

            # Image should occupy a meaningful portion of the slide
            # At minimum 20% of slide width and 15% of slide height
            width_ratio = pic.width / slide_width
            height_ratio = pic.height / slide_height

            # Image should be within slide bounds
            within_bounds = (pic.left >= 0 and pic.top >= 0 and
                             pic.left + pic.width <= slide_width + Inches(0.5) and
                             pic.top + pic.height <= slide_height + Inches(0.5))

            reasonable_size = width_ratio >= 0.2 and height_ratio >= 0.15

            if within_bounds and reasonable_size:
                print(f"PASS: Component 3 — Image positioned and sized appropriately "
                      f"(w_ratio={width_ratio:.2f}, h_ratio={height_ratio:.2f}) (0.2 pts)")
                total_score += 0.2
            elif within_bounds or reasonable_size:
                print(f"PARTIAL: Component 3 — Image partially OK "
                      f"(bounds={within_bounds}, size={reasonable_size}, "
                      f"w_ratio={width_ratio:.2f}, h_ratio={height_ratio:.2f}) (0.1 pts)")
                total_score += 0.1
            else:
                print(f"FAIL: Component 3 — Image too small or out of bounds "
                      f"(w_ratio={width_ratio:.2f}, h_ratio={height_ratio:.2f}, bounds={within_bounds})")
        else:
            print(f"FAIL: Component 3 — No picture to evaluate position/size")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
