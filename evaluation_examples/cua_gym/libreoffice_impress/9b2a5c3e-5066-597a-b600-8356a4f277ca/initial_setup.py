"""
Initial Setup: Broken OLE object on slide 6 of Data Review presentation
Task ID: impress_fix_053
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image, ImageDraw, ImageFont

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_053'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
DATA_TABLE_IMG = f'{WORKDIR}/Desktop/data_table.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_data_table_image():
    """Create a realistic data table screenshot image at 800x400."""
    os.makedirs(f'{WORKDIR}/Desktop', exist_ok=True)

    img = Image.new('RGB', (800, 400), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)

    # Try to use a decent font, fall back to default
    try:
        font_header = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 14)
        font_cell = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 12)
    except Exception:
        font_header = ImageFont.load_default()
        font_cell = ImageFont.load_default()

    # Table dimensions
    headers = ["Region", "Q1 Revenue", "Q2 Revenue", "Q3 Revenue", "Q4 Revenue", "Total"]
    data = [
        ["North America", "$1,245,800", "$1,389,200", "$1,512,400", "$1,678,900", "$5,826,300"],
        ["Europe", "$987,300", "$1,045,600", "$1,123,800", "$1,234,500", "$4,391,200"],
        ["Asia Pacific", "$756,200", "$834,100", "$912,500", "$1,023,400", "$3,526,200"],
        ["Latin America", "$345,600", "$378,900", "$412,300", "$456,700", "$1,593,500"],
        ["Middle East", "$234,100", "$256,800", "$278,400", "$312,600", "$1,081,900"],
        ["Africa", "$123,400", "$145,200", "$167,800", "$189,300", "$625,700"],
    ]

    col_widths = [120, 110, 110, 110, 110, 110]
    row_height = 35
    start_x, start_y = 30, 30

    # Title
    draw.text((start_x, 5), "Regional Revenue Summary - FY2025", fill=(33, 37, 41), font=font_header)

    # Draw header row
    x = start_x
    for i, header in enumerate(headers):
        draw.rectangle([x, start_y, x + col_widths[i], start_y + row_height],
                       fill=(52, 73, 94), outline=(44, 62, 80))
        draw.text((x + 8, start_y + 10), header, fill=(255, 255, 255), font=font_header)
        x += col_widths[i]

    # Draw data rows
    for r, row in enumerate(data):
        y = start_y + (r + 1) * row_height
        bg_color = (245, 247, 250) if r % 2 == 0 else (255, 255, 255)
        x = start_x
        for c, cell in enumerate(row):
            draw.rectangle([x, y, x + col_widths[c], y + row_height],
                           fill=bg_color, outline=(200, 200, 200))
            draw.text((x + 8, y + 10), cell, fill=(33, 37, 41), font=font_cell)
            x += col_widths[c]

    # Draw outer border
    total_w = sum(col_widths)
    total_h = (len(data) + 1) * row_height
    draw.rectangle([start_x, start_y, start_x + total_w, start_y + total_h],
                   outline=(44, 62, 80), width=2)

    img.save(DATA_TABLE_IMG)
    print(f'Data table image created: {DATA_TABLE_IMG}')


def add_text_to_slide(slide, text, left, top, width, height, font_size=18, bold=False,
                      color=RGBColor(0x21, 0x25, 0x29), alignment=PP_ALIGN.LEFT):
    """Helper to add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

    add_text_to_slide(slide1, "FY2025 Data Review", Inches(1.5), Inches(2),
                      Inches(10), Inches(1.5), font_size=40, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_text_to_slide(slide1, "Quarterly Performance Analysis\nPrepared by Finance Division",
                      Inches(2.5), Inches(4), Inches(8), Inches(1.5), font_size=20,
                      color=RGBColor(0xCC, 0xCC, 0xCC), alignment=PP_ALIGN.CENTER)

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_to_slide(slide2, "Agenda", Inches(0.8), Inches(0.4),
                      Inches(5), Inches(0.8), font_size=32, bold=True,
                      color=RGBColor(0x1B, 0x2A, 0x4A))

    agenda_items = [
        "1. Executive Summary",
        "2. Revenue Overview by Region",
        "3. Cost Analysis",
        "4. Profitability Metrics",
        "5. Regional Data Tables",
        "6. Key Takeaways & Next Steps",
    ]
    txBox = slide2.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(8), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(agenda_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.space_after = Pt(12)
        run = p.runs[0]
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 3: Executive Summary ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_to_slide(slide3, "Executive Summary", Inches(0.8), Inches(0.4),
                      Inches(6), Inches(0.8), font_size=32, bold=True,
                      color=RGBColor(0x1B, 0x2A, 0x4A))

    summary_text = (
        "Overall revenue grew 12.3% year-over-year, driven primarily by strong "
        "performance in North America and Asia Pacific regions. Operating margins "
        "improved by 2.1 percentage points to 18.7%, reflecting cost optimization "
        "initiatives launched in Q2. The EMEA region faced headwinds from currency "
        "fluctuations but showed resilient underlying growth of 8.5%."
    )
    add_text_to_slide(slide3, summary_text, Inches(0.8), Inches(1.5),
                      Inches(11), Inches(3), font_size=16,
                      color=RGBColor(0x44, 0x44, 0x44))

    # Key metrics boxes
    metrics = [("$17.0B", "Total Revenue"), ("18.7%", "Operating Margin"),
               ("+12.3%", "YoY Growth"), ("$3.2B", "Net Income")]
    for i, (value, label) in enumerate(metrics):
        left = Inches(0.8 + i * 3.1)
        shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        left, Inches(4.5), Inches(2.6), Inches(1.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
        shape.line.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = value
        run = p.runs[0]
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = label
        run2 = p2.runs[0]
        run2.font.size = Pt(14)
        run2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # --- Slide 4: Revenue by Region ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_to_slide(slide4, "Revenue by Region", Inches(0.8), Inches(0.4),
                      Inches(6), Inches(0.8), font_size=32, bold=True,
                      color=RGBColor(0x1B, 0x2A, 0x4A))

    add_text_to_slide(slide4, (
        "North America continues to be our largest market, contributing 34% of total "
        "revenue. Asia Pacific showed the highest growth rate at 18.2%, driven by expansion "
        "in Southeast Asian markets. Europe maintained steady performance despite macro "
        "challenges, while Latin America showed promising trajectory in H2."
    ), Inches(0.8), Inches(1.5), Inches(11), Inches(4.5), font_size=16,
       color=RGBColor(0x44, 0x44, 0x44))

    # --- Slide 5: Cost Analysis ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_to_slide(slide5, "Cost Analysis", Inches(0.8), Inches(0.4),
                      Inches(6), Inches(0.8), font_size=32, bold=True,
                      color=RGBColor(0x1B, 0x2A, 0x4A))

    # Add a small table
    table_shape = slide5.shapes.add_table(5, 3, Inches(1), Inches(1.5), Inches(9), Inches(3.5))
    table = table_shape.table
    table.columns[0].width = Inches(3)
    table.columns[1].width = Inches(3)
    table.columns[2].width = Inches(3)

    cost_data = [
        ["Cost Category", "FY2024", "FY2025"],
        ["Personnel", "$4,230,000", "$4,567,000"],
        ["Infrastructure", "$1,890,000", "$2,012,000"],
        ["Marketing", "$1,234,000", "$1,345,000"],
        ["R&D", "$2,567,000", "$2,890,000"],
    ]
    for r, row_data in enumerate(cost_data):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(14)
                if r == 0:
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Color header row
    for c in range(3):
        cell = table.cell(0, c)
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

    # --- Slide 6: Regional Data Tables (with broken OLE) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_to_slide(slide6, "Regional Data Tables", Inches(0.8), Inches(0.4),
                      Inches(6), Inches(0.8), font_size=32, bold=True,
                      color=RGBColor(0x1B, 0x2A, 0x4A))

    add_text_to_slide(slide6, "Source: Finance ERP System - Embedded Spreadsheet",
                      Inches(0.8), Inches(1.1), Inches(8), Inches(0.5), font_size=12,
                      color=RGBColor(0x88, 0x88, 0x88))

    # Create a "broken OLE object" - gray rectangle with X to simulate broken embed
    ole_left = Inches(1.5)
    ole_top = Inches(1.8)
    ole_width = Inches(8)
    ole_height = Inches(4.5)

    # Gray box shape simulating broken OLE object
    broken_ole = slide6.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, ole_left, ole_top, ole_width, ole_height
    )
    broken_ole.fill.solid()
    broken_ole.fill.fore_color.rgb = RGBColor(0xC0, 0xC0, 0xC0)  # Gray
    broken_ole.line.color.rgb = RGBColor(0x80, 0x80, 0x80)
    broken_ole.line.width = Pt(2)
    broken_ole.name = "Broken_OLE_Object"

    # Add X mark text inside or overlaid to indicate broken
    x_mark = slide6.shapes.add_textbox(
        ole_left + Inches(3), ole_top + Inches(1.2), Inches(2), Inches(2)
    )
    tf = x_mark.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "X"
    run = p.runs[0]
    run.font.size = Pt(72)
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    run.font.bold = True
    x_mark.name = "Broken_OLE_XMark"

    # Small label underneath
    label = slide6.shapes.add_textbox(
        ole_left + Inches(1.5), ole_top + Inches(3.5), Inches(5), Inches(0.5)
    )
    tf_label = label.text_frame
    p_label = tf_label.paragraphs[0]
    p_label.alignment = PP_ALIGN.CENTER
    p_label.text = "[ Embedded Object - Cannot be displayed ]"
    run_label = p_label.runs[0]
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
    run_label.font.italic = True
    label.name = "Broken_OLE_Label"

    # --- Slide 7: Profitability ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_to_slide(slide7, "Profitability Metrics", Inches(0.8), Inches(0.4),
                      Inches(6), Inches(0.8), font_size=32, bold=True,
                      color=RGBColor(0x1B, 0x2A, 0x4A))

    add_text_to_slide(slide7, (
        "Gross profit margin expanded to 42.3% in FY2025, up from 40.1% in FY2024. "
        "This improvement was primarily driven by operational efficiencies in the "
        "Asia Pacific region and favorable product mix shift toward higher-margin "
        "enterprise solutions. EBITDA margin reached 24.5%, exceeding target by 1.2pp."
    ), Inches(0.8), Inches(1.5), Inches(11), Inches(4), font_size=16,
       color=RGBColor(0x44, 0x44, 0x44))

    # --- Slide 8: Key Takeaways ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    fill8 = slide8.background.fill
    fill8.solid()
    fill8.fore_color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

    add_text_to_slide(slide8, "Key Takeaways & Next Steps", Inches(1), Inches(0.5),
                      Inches(10), Inches(1), font_size=36, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)

    takeaways = [
        "Revenue growth exceeded expectations across all major regions",
        "Cost optimization program delivered $180M in annual savings",
        "Strategic investments in APAC positioned for continued growth",
        "Q1 FY2026 guidance: 10-12% revenue growth target",
    ]
    txBox = slide8.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(takeaways):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"  {item}"
        p.space_after = Pt(16)
        run = p.runs[0]
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Create the data table image on Desktop
    create_data_table_image()

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
