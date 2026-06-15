"""
Initial Setup: Format text on slide 7 with multiple formatting levels
Task ID: impress_gf1_048
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf1_048'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Typography & Visual Design"
    slide1.placeholders[1].text = "A Guide to Effective Text Hierarchy"

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    tb = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Why Typography Matters"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.name = "Calibri"

    tb2 = slide2.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(4))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Good typography creates visual hierarchy, guiding the reader's eye through content in a deliberate sequence. Studies show that 95% of web design is typography, making it the single most important element in any visual communication."
    run2 = p2.runs[0]
    run2.font.size = Pt(18)
    run2.font.name = "Calibri"

    # --- Slide 3: Principles of Hierarchy ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Three Pillars of Typographic Hierarchy"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.name = "Calibri"

    tb2 = slide3.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    items = [
        "1. Size — Larger text draws attention first, establishing primary importance.",
        "2. Weight — Bold text carries more visual gravity within the same size tier.",
        "3. Color — Darker or more saturated tones create emphasis; lighter tones recede.",
    ]
    for i, item in enumerate(items):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = item
        run = p.runs[0]
        run.font.size = Pt(18)
        run.font.name = "Calibri"
        p.space_after = Pt(12)

    # --- Slide 4: Font Pairing ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide4.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Font Pairing Best Practices"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.name = "Calibri"

    tb2 = slide4.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "Contrast is key when pairing fonts. Combine a serif with a sans-serif, or a display font with a body font. Avoid combining fonts that are too similar — the subtle differences will look like a mistake rather than a design choice."
    run = p.runs[0]
    run.font.size = Pt(18)
    run.font.name = "Calibri"

    p2 = tf2.add_paragraph()
    p2.text = "Recommended pairs: Montserrat + Merriweather, Playfair Display + Source Sans Pro, Raleway + Lora."
    run2 = p2.runs[0]
    run2.font.size = Pt(16)
    run2.font.name = "Calibri"

    # --- Slide 5: Color in Typography ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Using Color to Create Depth"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.name = "Calibri"

    tb2 = slide5.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(4))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "A monochromatic palette with varying saturation and lightness levels creates an elegant depth effect. Dark tones anchor primary headings, mid-tones support secondary content, and light tones serve as subtle annotations."
    run = p.runs[0]
    run.font.size = Pt(18)
    run.font.name = "Calibri"

    # --- Slide 6: Spacing & Alignment ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide6.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Spacing and Alignment Techniques"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.name = "Calibri"

    tb2 = slide6.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(4))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "Generous line spacing (1.3x to 1.6x) improves readability. Consistent margins frame content and give it breathing room. Left alignment is safest for body text; center alignment works for titles and short blocks."
    run = p.runs[0]
    run.font.size = Pt(18)
    run.font.name = "Calibri"

    # --- Slide 7: THE KEY SLIDE — plain unformatted text, 3 paragraphs at 18pt ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide7.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(5))
    tf = tb.text_frame
    tf.word_wrap = True

    # Paragraph 1
    p1 = tf.paragraphs[0]
    p1.text = "The Core Message: Every great design begins with a clear typographic hierarchy that guides the reader effortlessly."
    run1 = p1.runs[0]
    run1.font.size = Pt(18)
    run1.font.name = "Calibri"
    # No bold, no italic, no color — plain default

    # Paragraph 2
    p2 = tf.add_paragraph()
    p2.text = "Supporting detail: When headings, subheadings, and body text each occupy a distinct visual tier, the audience absorbs information naturally without conscious effort."
    run2 = p2.runs[0]
    run2.font.size = Pt(18)
    run2.font.name = "Calibri"

    # Paragraph 3
    p3 = tf.add_paragraph()
    p3.text = 'Attribution: "Design is not just what it looks like and feels like. Design is how it works." — Steve Jobs'
    run3 = p3.runs[0]
    run3.font.size = Pt(18)
    run3.font.name = "Calibri"

    # --- Slide 8: Summary ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide8.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Takeaways"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.name = "Calibri"

    tb2 = slide8.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    items = [
        "Use size, weight, and color together to establish clear hierarchy levels.",
        "Limit your palette to 2-3 font families maximum for cohesion.",
        "Test your typography at multiple screen sizes to ensure readability.",
        "Remember: the goal of typography is communication, not decoration.",
    ]
    for i, item in enumerate(items):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = item
        run = p.runs[0]
        run.font.size = Pt(18)
        run.font.name = "Calibri"
        p.space_after = Pt(10)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
