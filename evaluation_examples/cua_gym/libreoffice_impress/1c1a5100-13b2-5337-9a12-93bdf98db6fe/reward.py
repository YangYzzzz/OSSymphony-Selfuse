"""
Reward Script: Format text on slide 7 with typographic hierarchy
Task ID: impress_gf1_048
Domain: libreoffice_impress
Scoring:
  Component 1 (0.35): Para 1 — Calibri 24pt bold, color #0D1B2A
  Component 2 (0.35): Para 2 — Calibri 18pt regular, color #546E7A
  Component 3 (0.30): Para 3 — Calibri 14pt italic, color #90A4AE
"""

import os

from pptx import Presentation
from pptx.util import Pt

WORKDIR = '/home/user'
TASK_ID = 'impress_gf1_048'


def normalize_bool(val):
    """Normalize None/False to False, True to True for bold/italic."""
    return True if val is True else False


def get_run_color_hex(run):
    """Safely get the RGB hex string of a run's font color, or None."""
    try:
        if run.font.color.type is not None:
            return str(run.font.color.rgb).upper()
    except Exception:
        pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 7 slides
    if len(prs.slides) < 7:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 7")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[6]  # 0-indexed, slide 7

    # Find the text box with the 3 paragraphs (skip empty placeholders)
    target_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame:
            nonempty = [p for p in shape.text_frame.paragraphs if p.text.strip()]
            if len(nonempty) >= 3:
                # Check that paragraphs start with the expected prefixes
                texts = [p.text.strip() for p in nonempty]
                if texts[0].startswith("The Core Message") and texts[1].startswith("Supporting detail") and texts[2].startswith("Attribution"):
                    target_shape = shape
                    break

    if target_shape is None:
        print("FAIL: Could not find target text box on slide 7 with 3 expected paragraphs")
        print("REWARD: 0.0")
        return 0.0

    # Get the three non-empty paragraphs
    nonempty_paras = [p for p in target_shape.text_frame.paragraphs if p.text.strip()]
    para1, para2, para3 = nonempty_paras[0], nonempty_paras[1], nonempty_paras[2]

    # Component 1: Paragraph 1 — Calibri 24pt bold, color #0D1B2A (0.35 points)
    # Checks: size=24pt, bold=True, color=#0D1B2A (all differ from initial 18pt/None/None)
    try:
        runs1 = [r for r in para1.runs if (r.text or "").strip()]
        if not runs1:
            print("FAIL: Component 1 — Para 1 has no non-empty runs")
        else:
            sub_checks = 0
            run = runs1[0]

            # Check font size = 24pt (304800 EMU)
            if run.font.size is not None and run.font.size == Pt(24):
                sub_checks += 1
                print(f"PASS: Component 1a — Para 1 font size is 24pt ({run.font.size})")
            else:
                print(f"FAIL: Component 1a — Para 1 font size expected 24pt (304800), got {run.font.size}")

            # Check bold = True
            if normalize_bool(run.font.bold) is True:
                sub_checks += 1
                print(f"PASS: Component 1b — Para 1 is bold")
            else:
                print(f"FAIL: Component 1b — Para 1 bold expected True, got {run.font.bold}")

            # Check color = #0D1B2A
            color_hex = get_run_color_hex(run)
            if color_hex == "0D1B2A":
                sub_checks += 1
                print(f"PASS: Component 1c — Para 1 color is #0D1B2A")
            else:
                print(f"FAIL: Component 1c — Para 1 color expected #0D1B2A, got {color_hex}")

            # Award proportional credit: all 3 sub-checks must pass for full points
            if sub_checks == 3:
                total_score += 0.35
                print(f"PASS: Component 1 — Para 1 fully correct (0.35 pts)")
            elif sub_checks > 0:
                partial = round(0.35 * sub_checks / 3, 3)
                total_score += partial
                print(f"PARTIAL: Component 1 — Para 1 {sub_checks}/3 checks passed ({partial} pts)")
            else:
                print(f"FAIL: Component 1 — Para 1 no checks passed (0.0 pts)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Paragraph 2 — Calibri 18pt regular (not bold/italic), color #546E7A (0.35 points)
    # In initial state: 18pt, no bold, no italic, no color — the DIFFERENTIATOR is color #546E7A
    # Also explicitly set bold=False and italic=False (vs None in initial)
    try:
        runs2 = [r for r in para2.runs if (r.text or "").strip()]
        if not runs2:
            print("FAIL: Component 2 — Para 2 has no non-empty runs")
        else:
            sub_checks = 0
            run = runs2[0]

            # Check font size = 18pt (228600 EMU)
            if run.font.size is not None and run.font.size == Pt(18):
                sub_checks += 1
                print(f"PASS: Component 2a — Para 2 font size is 18pt ({run.font.size})")
            else:
                print(f"FAIL: Component 2a — Para 2 font size expected 18pt (228600), got {run.font.size}")

            # Check not bold and not italic
            if normalize_bool(run.font.bold) is False and normalize_bool(run.font.italic) is False:
                sub_checks += 1
                print(f"PASS: Component 2b — Para 2 is regular (not bold, not italic)")
            else:
                print(f"FAIL: Component 2b — Para 2 expected regular, got bold={run.font.bold} italic={run.font.italic}")

            # Check color = #546E7A (THIS is the key differentiator from initial state)
            color_hex = get_run_color_hex(run)
            if color_hex == "546E7A":
                sub_checks += 1
                print(f"PASS: Component 2c — Para 2 color is #546E7A")
            else:
                print(f"FAIL: Component 2c — Para 2 color expected #546E7A, got {color_hex}")

            # Only award points if color check passes (the task-introduced change)
            if sub_checks == 3:
                total_score += 0.35
                print(f"PASS: Component 2 — Para 2 fully correct (0.35 pts)")
            elif color_hex == "546E7A" and sub_checks >= 2:
                # Color is correct, one minor issue
                partial = round(0.35 * sub_checks / 3, 3)
                total_score += partial
                print(f"PARTIAL: Component 2 — Para 2 {sub_checks}/3 checks passed ({partial} pts)")
            elif color_hex == "546E7A":
                total_score += 0.117
                print(f"PARTIAL: Component 2 — Only color correct (0.117 pts)")
            else:
                print(f"FAIL: Component 2 — Color #546E7A not set (0.0 pts)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Paragraph 3 — Calibri 14pt italic, color #90A4AE (0.30 points)
    # Checks: size=14pt, italic=True, color=#90A4AE (all differ from initial 18pt/None/None)
    try:
        runs3 = [r for r in para3.runs if (r.text or "").strip()]
        if not runs3:
            print("FAIL: Component 3 — Para 3 has no non-empty runs")
        else:
            sub_checks = 0
            run = runs3[0]

            # Check font size = 14pt (177800 EMU)
            if run.font.size is not None and run.font.size == Pt(14):
                sub_checks += 1
                print(f"PASS: Component 3a — Para 3 font size is 14pt ({run.font.size})")
            else:
                print(f"FAIL: Component 3a — Para 3 font size expected 14pt (177800), got {run.font.size}")

            # Check italic = True
            if normalize_bool(run.font.italic) is True:
                sub_checks += 1
                print(f"PASS: Component 3b — Para 3 is italic")
            else:
                print(f"FAIL: Component 3b — Para 3 italic expected True, got {run.font.italic}")

            # Check color = #90A4AE
            color_hex = get_run_color_hex(run)
            if color_hex == "90A4AE":
                sub_checks += 1
                print(f"PASS: Component 3c — Para 3 color is #90A4AE")
            else:
                print(f"FAIL: Component 3c — Para 3 color expected #90A4AE, got {color_hex}")

            # Award proportional credit
            if sub_checks == 3:
                total_score += 0.30
                print(f"PASS: Component 3 — Para 3 fully correct (0.30 pts)")
            elif sub_checks > 0:
                partial = round(0.30 * sub_checks / 3, 3)
                total_score += partial
                print(f"PARTIAL: Component 3 — Para 3 {sub_checks}/3 checks passed ({partial} pts)")
            else:
                print(f"FAIL: Component 3 — Para 3 no checks passed (0.0 pts)")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook: save any unsaved LibreOffice edits before verifying
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
