"""
FINAL REWARD SCRIPT - SUCCESS
Task: In LibreOffice Impress, please open slide 35 and update all body text so it uses the Noto Sans font at exactly 20 pt, with the font color set to Dark Gray 30 % from the standard palette.
Generated: 2025-09-10 23:43:29
Status: success
Model: azure-o3
Total Steps: 5
"""

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Pt
import os


def verify_body_text_formatting(file_path: str) -> float:
    """Verify slide 35 body text formatting.

    Requirements (LibreOffice Impress task):
      1. Slide 35 must exist.
      2. All BODY text (exclude title/subtitle placeholders) on that slide
         • uses Noto Sans font
         • is exactly 20 pt
         • has colour Dark Gray 30 % (hex 4C4C4C)

    Progressive scoring (total 1.0):
      • 0.4 – slide 35 present
      • 0.3 – all body text font correct
      • 0.2 – all body text size correct
      • 0.1 – all body text colour correct
    """

    max_score = 1.0
    score = 0.0
    print(f"Verifying presentation: {file_path}")

    # ---------- File existence ----------
    if not os.path.exists(file_path):
        print("✗ File not found")
        return 0.0  # No progress possible

    # ---------- Load presentation ----------
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Failed to open presentation: {e}")
        return 0.0

    # ---------- Requirement 1: Slide 35 exists ----------
    slide_index = 34  # zero-based index
    if len(prs.slides) <= slide_index:
        print(f"✗ Slide 35 not found (presentation has {len(prs.slides)} slides)")
        return 0.0
    print("✓ Slide 35 found")
    score += 0.4

    slide = prs.slides[slide_index]

    # ---------- Gather body-text runs (exclude title/subtitle) ----------
    body_runs = []
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or shape.text_frame is None:
            continue  # skip non-text shapes

        # Skip title-type placeholders (we only care about body text)
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            if ph_type in (
                PP_PLACEHOLDER.TITLE,
                PP_PLACEHOLDER.CENTER_TITLE,
                PP_PLACEHOLDER.SUBTITLE,
            ):
                continue

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text and run.text.strip():  # ignore empty runs
                    body_runs.append(run)

    if not body_runs:
        print("✗ No body-text runs found on slide 35")
        return score  # Return partial credit (only slide existence)

    print(f"Found {len(body_runs)} body-text runs to evaluate")

    # ---------- Target formatting ----------
    target_font_names = {"Noto Sans", "NotoSans", "Noto Sans Regular"}
    target_size_emu = Pt(20)  # pptx util Pt returns EMU length object
    target_hex = "4C4C4C"  # Dark Gray 30 %

    font_ok = True
    size_ok = True
    color_ok = True

    for run in body_runs:
        font = run.font

        # --- Font name check ---
        if not (font.name and font.name.strip() in target_font_names):
            font_ok = False

        # --- Font size check ---
        if font.size != target_size_emu:
            size_ok = False

        # --- Font colour check ---
        this_color_ok = False
        if font.color is not None and font.color.type == 1 and font.color.rgb is not None:
            if str(font.color.rgb).upper() == target_hex:
                this_color_ok = True
        if not this_color_ok:
            color_ok = False

    # ---------- Scoring for each formatting criterion ----------
    if font_ok:
        print("✓ All body text uses Noto Sans font")
        score += 0.3
    else:
        print("✗ Some body text does not use Noto Sans font")

    if size_ok:
        print("✓ All body text is exactly 20 pt")
        score += 0.2
    else:
        print("✗ Some body text is not 20 pt")

    if color_ok:
        print("✓ All body text is Dark Gray 30 % (hex 4C4C4C)")
        score += 0.1
    else:
        print("✗ Some body text does not have the correct colour")

    # ---------- Final score (cap at 1.0 & fix float precision) ----------
    final_score = min(score, max_score)
    if final_score > 0.999:
        final_score = 1.0

    print(f"Total Score: {final_score}")
    return final_score


if __name__ == "__main__":
    FILE_PATH = "/home/user/in_libreoffice_impress_please_open_slide_35_and_update_all_body_text_so_it_uses_the_noto_sans_font_a_golden.pptx"
    reward = verify_body_text_formatting(FILE_PATH)
    print(f"REWARD: {reward}")

