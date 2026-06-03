"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 27 is the odd one out in my deck—everything else is dark, but I need that particular slide to sit on a lighter backdrop so the photos pop. In LibreOffice Impress, how do I set ONLY slide 27’s background to the palette shade “Light Gray 2” (#D9D9D9) without changing any of the other slides?
Generated: 2025-09-10 23:12:05
Status: success
Model: azure-o3
Total Steps: 2
"""

from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
import os

"""
Reward Script for LibreOffice Impress Task
Task: Ensure ONLY slide 27 has a solid background set to the palette shade
      “Light Gray 2” (#D9D9D9) while every other slide keeps a different background.

Scoring Breakdown (progressive):
  0.1  – Deck contains at least 27 slides (prerequisite for task relevance)
  0.1  – Slide 27 background fill is SOLID (not none/gradient/picture/etc.)
  0.6  – Slide 27 background colour is exactly #D9D9D9 (Light Gray 2)
  0.2  – No other slide (except 27) uses #D9D9D9 as its solid background colour
TOTAL = 1.0 for perfect completion

The script prints detailed verification steps and the final reward in the
required format:  "REWARD: X.X"
"""

def verify_slide27_background(file_path: str) -> float:
    """Verify that only slide 27 has Light Gray 2 background (#D9D9D9)."""
    print(f"Verifying presentation file: {file_path}")
    total_score = 0.0
    max_score = 1.0

    # 1. Load the presentation ------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File does not exist")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        return 0.0

    slide_count = len(prs.slides)
    print(f"Presentation contains {slide_count} slides")

    # 2. Check slide count (≥27) ---------------------------------------------
    if slide_count >= 27:
        total_score += 0.1
        print("✓ At least 27 slides present (0.1)")
    else:
        print("✗ Less than 27 slides – task context invalid")
        print(f"REWARD: {total_score}")
        return total_score  # nothing more to verify

    # 3. Inspect slide 27 background -----------------------------------------
    slide_27 = prs.slides[26]  # zero-based index
    fill = slide_27.background.fill

    if fill.type == MSO_FILL_TYPE.SOLID:
        total_score += 0.1
        print("✓ Slide 27 has SOLID fill (0.1)")
        try:
            color_hex = str(fill.fore_color.rgb).upper()
        except Exception:
            color_hex = None
        print(f"Slide 27 colour detected: {color_hex}")

        if color_hex == "D9D9D9":
            total_score += 0.6
            print("✓ Slide 27 colour is Light Gray 2 (#D9D9D9) (0.6)")
        else:
            print("✗ Slide 27 colour is not Light Gray 2 – no points for colour match")
    else:
        print("✗ Slide 27 background is not solid – no points for colour checks")

    # 4. Ensure exclusivity of colour on other slides -------------------------
    other_matches = 0
    for idx, slide in enumerate(prs.slides):
        if idx == 26:
            continue  # skip slide 27 itself
        f = slide.background.fill
        if f.type == MSO_FILL_TYPE.SOLID:
            try:
                if str(f.fore_color.rgb).upper() == "D9D9D9":
                    other_matches += 1
            except Exception:
                # ignore shapes without colour information
                pass

    if other_matches == 0:
        total_score += 0.2
        print("✓ No other slides have Light Gray 2 background (0.2)")
    else:
        print(f"✗ {other_matches} other slide(s) share Light Gray 2 background – exclusivity failed")

    # 5. Final score ----------------------------------------------------------
    final_score = min(total_score, max_score)
    print(f"Total score: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score


# -----------------------------
# Execute verification function
# -----------------------------
file_path = "/home/user/slide_27_is_the_odd_one_out_in_my_deckeverything_else_is_dark_but_i_need_that_particular_slide_to_si_golden.pptx"
verify_slide27_background(file_path)
