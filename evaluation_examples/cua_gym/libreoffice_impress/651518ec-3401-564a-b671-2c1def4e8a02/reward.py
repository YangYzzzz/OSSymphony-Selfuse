"""
Reward Script: Change all text throughout presentation to Palatino Linotype 14pt
Task ID: osworld_impress_global_font_change_012
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): All runs across all slides use 'Palatino Linotype' font
  Component 2 (0.5): All runs across all slides use 14pt (177800 EMU) font size
"""

import os

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_global_font_change_012'

TARGET_FONT = 'Palatino Linotype'
TARGET_SIZE_EMU = 177800  # 14pt * 12700 = 177800 EMU


def get_all_runs(prs):
    """Collect all non-empty runs from all text frames on all slides."""
    runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if (run.text or "").strip():
                            runs.append(run)
    return runs


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition gate: verify 8 slides exist
    num_slides = len(prs.slides)
    if num_slides != 8:
        print(f"CRITICAL: Expected 8 slides, found {num_slides}. File may be corrupt.")
        print("REWARD: 0.0")
        return 0.0

    # Collect all non-empty runs
    try:
        all_runs = get_all_runs(prs)
    except Exception as e:
        print(f"CRITICAL: Cannot enumerate runs: {e}")
        print("REWARD: 0.0")
        return 0.0

    if not all_runs:
        print("CRITICAL: No non-empty runs found in presentation.")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: All runs use 'Palatino Linotype' font (0.5 points)
    # The task requires all text in all textboxes on all 8 slides to use Palatino Linotype.
    # In the initial file, fonts include Arial, Calibri, Georgia, Verdana, Times New Roman.
    # A correct completion changes every run to Palatino Linotype.
    try:
        wrong_font_runs = []
        for run in all_runs:
            actual_font = run.font.name
            if actual_font != TARGET_FONT:
                wrong_font_runs.append((run.text[:30], actual_font))

        if not wrong_font_runs:
            print(f"PASS: Component 1 — All {len(all_runs)} runs use '{TARGET_FONT}' font (0.5 pts)")
            total_score += 0.5
        else:
            examples = wrong_font_runs[:5]
            print(f"FAIL: Component 1 — {len(wrong_font_runs)}/{len(all_runs)} runs do NOT use '{TARGET_FONT}'")
            for text, font in examples:
                print(f"  run '{text}' has font='{font}'")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: All runs use 14pt (177800 EMU) font size (0.5 points)
    # The task requires all text to be 14pt. In the initial file, sizes vary widely:
    # titles use 355600-457200 EMU, body text 177800-203200, footnotes 127000-139700.
    # A correct completion changes all run sizes to exactly 177800 (14pt).
    try:
        wrong_size_runs = []
        for run in all_runs:
            actual_size = run.font.size
            if actual_size != TARGET_SIZE_EMU:
                wrong_size_runs.append((run.text[:30], actual_size))

        if not wrong_size_runs:
            print(f"PASS: Component 2 — All {len(all_runs)} runs use 14pt ({TARGET_SIZE_EMU} EMU) size (0.5 pts)")
            total_score += 0.5
        else:
            examples = wrong_size_runs[:5]
            print(f"FAIL: Component 2 — {len(wrong_size_runs)}/{len(all_runs)} runs do NOT use {TARGET_SIZE_EMU} EMU")
            for text, size in examples:
                print(f"  run '{text}' has size={size}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
