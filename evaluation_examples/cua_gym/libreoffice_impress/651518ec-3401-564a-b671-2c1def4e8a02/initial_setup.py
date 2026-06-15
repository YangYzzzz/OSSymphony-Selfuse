"""
Initial Setup: 8-slide academic presentation on climate science with varied fonts/sizes
Task ID: osworld_impress_global_font_change_012
Domain: libreoffice_impress

NOTE: Intentionally uses mixed fonts (Arial, Calibri, Times New Roman, Verdana, Georgia)
and varied sizes — NOT Palatino Linotype and NOT uniform 14pt.
The task is to change ALL text to Palatino Linotype at 14pt.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_global_font_change_012'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_text_run(run, text, font_name, font_size_pt, bold=False, italic=False, color=None):
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = RGBColor(*color)


def add_text_box(slide, left_in, top_in, width_in, height_in,
                 lines, font_name, font_size_pt, bold=False, italic=False, color=None):
    """Add a text box with multiple lines, each as a paragraph."""
    txBox = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        run = para.add_run()
        set_text_run(run, line, font_name, font_size_pt, bold=bold, italic=italic, color=color)
    return txBox


def create_initial():
    prs = Presentation()
    # Standard widescreen: 10 x 7.5 inches
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title Slide ──────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
    title1 = slide1.shapes.title
    title1.text = "Climate Science: Understanding Our Changing Planet"
    for para in title1.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Arial"
            run.font.size = Pt(36)
            run.font.bold = True

    subtitle1 = slide1.placeholders[1]
    subtitle1.text = "A Comprehensive Academic Overview"
    for para in subtitle1.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(20)
            run.font.italic = True

    # Footer / author on slide 1
    add_text_box(slide1, 0.5, 6.5, 9, 0.6,
                 ["Dr. Elena Vasquez | Department of Atmospheric Sciences | Spring 2025"],
                 "Georgia", 11, italic=True, color=(80, 80, 80))

    # ── Slide 2: Introduction ─────────────────────────────────────────────────
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide2.shapes.title.text = "Introduction to Climate Science"
    for para in slide2.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Arial"
            run.font.size = Pt(32)
            run.font.bold = True

    body2 = slide2.placeholders[1]
    body2.text = ""
    bullets2 = [
        "Climate science studies long-term patterns in Earth's atmosphere",
        "Key variables: temperature, precipitation, wind patterns, ocean currents",
        "Evidence drawn from ice cores, tree rings, satellite data, and ocean sediments",
        "Distinction between weather (days) and climate (decades to centuries)",
        "Interdisciplinary field combining physics, chemistry, biology, and oceanography",
    ]
    tf2 = body2.text_frame
    for i, b in enumerate(bullets2):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = b
        run.font.name = "Verdana"
        run.font.size = Pt(16)

    add_text_box(slide2, 0.3, 6.8, 9.4, 0.4,
                 ["Source: IPCC Sixth Assessment Report, 2021"],
                 "Times New Roman", 10, italic=True, color=(100, 100, 100))

    # ── Slide 3: CO2 Levels ───────────────────────────────────────────────────
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Atmospheric CO₂ Concentrations"
    for para in slide3.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(30)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    body3 = slide3.placeholders[1]
    body3.text = ""
    bullets3 = [
        "Pre-industrial CO₂: ~280 ppm (parts per million)",
        "2023 annual average: 421.08 ppm — highest in 800,000 years",
        "Rate of increase: ~2.5 ppm/year (accelerating since 1960s)",
        "Primary sources: fossil fuel combustion, deforestation, cement production",
        "Carbon sinks: oceans (~26%), terrestrial vegetation (~30%)",
        "Keeling Curve documents continuous rise since 1958 (Mauna Loa Observatory)",
    ]
    tf3 = body3.text_frame
    for i, b in enumerate(bullets3):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = b
        run.font.name = "Arial"
        run.font.size = Pt(15)

    add_text_box(slide3, 0.5, 6.3, 9, 0.5,
                 ["Fig. 1 — Monthly mean CO₂ concentration (ppm), 1958–2024, NOAA Global Monitoring Laboratory"],
                 "Georgia", 10, italic=True)

    # ── Slide 4: Temperature Trends ───────────────────────────────────────────
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Global Mean Surface Temperature Trends"
    for para in slide4.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Times New Roman"
            run.font.size = Pt(28)
            run.font.bold = True

    body4 = slide4.placeholders[1]
    body4.text = ""
    bullets4 = [
        "Global average temperature has risen ~1.1 °C above pre-industrial baseline (1850–1900)",
        "The last decade (2011–2020) was the warmest on record",
        "Hottest individual year to date: 2023 (+1.45 °C anomaly, Copernicus/ECMWF)",
        "Polar amplification: Arctic warming 3–4× faster than global average",
        "Urban heat island effect complicates station-based measurements",
        "Homogenization algorithms correct for station moves, equipment changes",
    ]
    tf4 = body4.text_frame
    for i, b in enumerate(bullets4):
        if i == 0:
            p = tf4.paragraphs[0]
        else:
            p = tf4.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = b
        run.font.name = "Verdana"
        run.font.size = Pt(14)

    # Sub-caption box
    add_text_box(slide4, 0.5, 6.2, 9, 0.7,
                 ["Data: HadCRUT5, GISTEMP v4, Berkeley Earth Surface Temperature (BEST)",
                  "Note: anomalies calculated relative to 1951–1980 baseline"],
                 "Calibri", 10, italic=True, color=(90, 90, 90))

    # ── Slide 5: Sea Level Rise ───────────────────────────────────────────────
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Sea Level Rise: Observations and Projections"
    for para in slide5.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x00, 0x5B, 0x96)

    body5 = slide5.placeholders[1]
    body5.text = ""
    bullets5 = [
        "Global mean sea level has risen ~20 cm since 1901",
        "Rate accelerating: ~3.7 mm/year in 2006–2018 vs ~1.4 mm/year in 1901–1990",
        "Contributors: thermal expansion (~40%), glacier melt (~25%), ice sheets (~35%)",
        "Greenland ice sheet losing ~280 Gt/year; West Antarctica ~150 Gt/year",
        "Projections: 0.28–1.01 m rise by 2100 (SSP1-2.6 to SSP5-8.5 scenarios)",
        "Coastal flooding impacts ~1 billion people in low-elevation coastal zones",
    ]
    tf5 = body5.text_frame
    for i, b in enumerate(bullets5):
        if i == 0:
            p = tf5.paragraphs[0]
        else:
            p = tf5.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = b
        run.font.name = "Arial"
        run.font.size = Pt(15)

    add_text_box(slide5, 0.5, 6.3, 9, 0.5,
                 ["Source: IPCC AR6, Chapter 9; NASA Goddard Space Flight Center altimetry data"],
                 "Times New Roman", 10, italic=True, color=(100, 100, 100))

    # ── Slide 6: Arctic Sea Ice ───────────────────────────────────────────────
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Arctic Sea Ice Decline"
    for para in slide6.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Georgia"
            run.font.size = Pt(32)
            run.font.bold = True

    body6 = slide6.placeholders[1]
    body6.text = ""
    bullets6 = [
        "Arctic minimum sea ice extent declined ~13% per decade since 1979",
        "Record low September extent: 3.41 million km² (September 2012)",
        "Multi-year ice (≥5 years old) declined >90% since 1985",
        "Ice-albedo feedback: loss of reflective ice → ocean absorbs more solar radiation",
        "Ice-free Arctic summers possible before 2050 under high-emissions scenarios",
        "Consequences: altered jet stream, methane release from permafrost, ecosystem shifts",
    ]
    tf6 = body6.text_frame
    for i, b in enumerate(bullets6):
        if i == 0:
            p = tf6.paragraphs[0]
        else:
            p = tf6.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = b
        run.font.name = "Verdana"
        run.font.size = Pt(14)

    add_text_box(slide6, 0.5, 6.3, 9, 0.5,
                 ["Fig. 4 — September Arctic sea ice extent 1979–2024 (NSIDC / JAXA AMSR2)"],
                 "Calibri", 10, italic=True, color=(80, 80, 80))

    # ── Slide 7: Policy Implications ─────────────────────────────────────────
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Policy Implications and Mitigation Pathways"
    for para in slide7.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Arial"
            run.font.size = Pt(26)
            run.font.bold = True

    body7 = slide7.placeholders[1]
    body7.text = ""
    bullets7 = [
        "Paris Agreement (2015): limit warming to well below 2 °C, pursue 1.5 °C",
        "Net-zero emissions by 2050 required to meet 1.5 °C target (IPCC SR1.5)",
        "Mitigation: rapid decarbonization of energy, transport, industry, and land use",
        "Carbon pricing mechanisms: emissions trading systems (ETS) and carbon taxes",
        "Negative emissions technologies: BECCS, direct air capture, enhanced weathering",
        "Adaptation: coastal defenses, drought-resistant crops, early warning systems",
        "Climate finance: $100B/year commitment to developing nations (Green Climate Fund)",
    ]
    tf7 = body7.text_frame
    for i, b in enumerate(bullets7):
        if i == 0:
            p = tf7.paragraphs[0]
        else:
            p = tf7.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = b
        run.font.name = "Times New Roman"
        run.font.size = Pt(14)

    add_text_box(slide7, 0.5, 6.5, 9, 0.6,
                 ["Reference: UNEP Emissions Gap Report 2023; IPCC AR6 Working Group III"],
                 "Georgia", 10, italic=True, color=(90, 90, 90))

    # ── Slide 8: Conclusion ───────────────────────────────────────────────────
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Conclusions and Future Directions"
    for para in slide8.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(30)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    body8 = slide8.placeholders[1]
    body8.text = ""
    bullets8 = [
        "Human-caused climate change is unequivocal (IPCC AR6 consensus)",
        "All major indicators are trending in the same direction: warming",
        "Urgent, transformative action is necessary across all sectors",
        "Science provides the roadmap; implementation requires political will and equity",
        "Emerging research areas: tipping points, climate attribution, solar geoengineering",
        "The window for limiting warming to 1.5 °C is rapidly closing",
    ]
    tf8 = body8.text_frame
    for i, b in enumerate(bullets8):
        if i == 0:
            p = tf8.paragraphs[0]
        else:
            p = tf8.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = b
        run.font.name = "Arial"
        run.font.size = Pt(16)

    add_text_box(slide8, 0.5, 6.2, 9, 0.7,
                 ["\"The evidence is clear: climate change is a threat to human well-being and the health of the planet.\"",
                  "— IPCC AR6 Synthesis Report, 2023"],
                 "Georgia", 11, italic=True, color=(60, 60, 60))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
