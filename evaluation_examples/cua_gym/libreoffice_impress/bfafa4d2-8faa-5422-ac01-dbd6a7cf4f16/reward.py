"""
Reward Script: Copy talking points from pitch_talking_points.docx into
               the notes of slides 1-7 in Startup_Pitch.pptx.
Task ID: osworld_multi_apps_impress_notes_import_005
Domain: libreoffice_impress
Scoring:
  Component 1: All 7 slides have non-empty notes            — 0.30 points
  Component 2: Slide 1 notes contain expected bullet text   — 0.20 points
  Component 3: Slides 2-5 notes contain expected text       — 0.30 points
  Component 4: Slides 6-7 notes contain expected text       — 0.20 points
  Total: 1.00
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_multi_apps_impress_notes_import_005'

# Expected notes text for each slide (key phrases from pitch_talking_points.docx).
EXPECTED_NOTES = {
    1: [
        "Welcome the audience and introduce yourself as CEO and co-founder",
        "NovaMind AI was born out of personal frustration automating enterprise workflows",
        "We are here today because we believe AI-driven automation should be accessible to every business",
        "Our team has 40+ years of combined experience in enterprise software and AI",
    ],
    2: [
        "Every day, knowledge workers spend nearly half their time on tasks that could be automated",
        "Legacy RPA tools are brittle, expensive, and require dedicated engineering teams to maintain",
        "Small and mid-sized businesses are priced out of enterprise automation entirely",
        "The result: billions of dollars in lost productivity and missed revenue opportunities annually",
        "We validated this with 200+ interviews across finance, healthcare, and logistics sectors",
    ],
    3: [
        "NovaMind AI lets any business user describe a workflow in plain English and deploy it instantly",
        "Our AI engine maps intent to actions across connected apps with no coding required",
        "Security-first architecture: SOC 2 Type II certified, data never leaves your cloud environment",
        "We integrate with the tools your teams already use",
    ],
    4: [
        "The intelligent process automation market is growing faster than analysts predicted",
        "Regulation changes in finance and healthcare are accelerating digital transformation",
        "Our sweet spot: companies with 50",
        "Competitive moat: proprietary workflow graph model trained on 10M+ enterprise tasks",
    ],
    5: [
        "We are not pre-revenue",
        "Meridian Health reduced claims processing time by 67% in 90 days",
        "TerraLogix automated their entire vendor onboarding process, saving 1,200 hours per quarter",
        "Our NPS score is 72",
        "Pipeline: 45 active enterprise pilots expected to convert in Q2 2026",
    ],
    6: [
        "Simple, predictable pricing aligned with customer value delivery",
        "We expand naturally as customers automate more workflows",
        "Customer acquisition cost: $8,400",
        "Unit economics strengthen at scale as AI model accuracy improves with more data",
    ],
    7: [
        "We are raising $12M to accelerate what is already working",
        "Series A will fund expansion into the UK, Germany, and Australia",
        "AI Copilot feature launching Q4 2026",
        "We are in conversations with three strategic VCs",
        "Thank you",
    ],
}


def get_slide_notes_text(slide):
    """Return the notes text for a slide, stripped."""
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def count_matching_bullets(notes_text, expected_bullets):
    """Return count of expected_bullets found in notes_text."""
    return sum(1 for b in expected_bullets if b in notes_text)


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have exactly 7 slides
    if len(prs.slides) != 7:
        print(f"FAIL: Expected 7 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Collect all notes
    all_notes = {}
    for i, slide in enumerate(prs.slides, start=1):
        all_notes[i] = get_slide_notes_text(slide)

    # ----------------------------------------------------------------
    # Component 1: All 7 slides have non-empty notes (0.30 points)
    # This FAILS on initial_env (all notes empty) and PASSES on golden_env.
    # ----------------------------------------------------------------
    try:
        slides_with_notes = sum(1 for i in range(1, 8) if all_notes[i].strip())
        if slides_with_notes == 7:
            print(f"PASS: Component 1 — All 7 slides have non-empty notes (0.30 pts)")
            total_score += 0.30
        else:
            empty_slides = [i for i in range(1, 8) if not all_notes[i].strip()]
            print(f"FAIL: Component 1 — Slides {empty_slides} still empty. Only {slides_with_notes}/7 populated.")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # ----------------------------------------------------------------
    # Component 2: Slide 1 notes contain expected bullet points (0.20 points)
    # Per bullet: 0.20/4 = 0.05 points each
    # This FAILS on initial_env (empty notes) and PASSES on golden_env.
    # ----------------------------------------------------------------
    try:
        slide1_notes = all_notes[1]
        slide1_bullets = EXPECTED_NOTES[1]
        n_found = count_matching_bullets(slide1_notes, slide1_bullets)
        n_total = len(slide1_bullets)
        per_bullet_pts = 0.20 / n_total
        if n_found == n_total:
            print(f"PASS: Component 2 — Slide 1: all {n_total} bullets found (0.20 pts)")
            total_score += 0.20
        elif n_found > 0:
            partial_pts = round(n_found * per_bullet_pts, 4)
            print(f"PARTIAL: Component 2 — Slide 1: {n_found}/{n_total} bullets ({partial_pts} pts)")
            total_score += partial_pts
        else:
            print(f"FAIL: Component 2 — Slide 1 notes empty or missing bullets.")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # ----------------------------------------------------------------
    # Component 3: Slides 2-5 notes contain expected text (0.30 points)
    # Per slide: 0.075 points, each bullet within that slide is proportional.
    # This FAILS on initial_env (empty notes) and PASSES on golden_env.
    # ----------------------------------------------------------------
    try:
        comp3_score = 0.0
        for slide_num in range(2, 6):
            notes_text = all_notes[slide_num]
            expected_bullets = EXPECTED_NOTES[slide_num]
            n_found = count_matching_bullets(notes_text, expected_bullets)
            n_total = len(expected_bullets)
            per_bullet_pts = 0.075 / n_total
            if n_found == n_total:
                print(f"PASS: Component 3 — Slide {slide_num}: all {n_total} bullets found")
                comp3_score += 0.075
            elif n_found > 0:
                slide_pts = round(n_found * per_bullet_pts, 4)
                print(f"PARTIAL: Component 3 — Slide {slide_num}: {n_found}/{n_total} bullets ({slide_pts} pts)")
                comp3_score += slide_pts
            else:
                print(f"FAIL: Component 3 — Slide {slide_num}: 0/{n_total} bullets. Notes: {repr(notes_text[:60])}")
        comp3_score = round(min(comp3_score, 0.30), 4)
        print(f"Component 3 subtotal: {comp3_score}/0.30")
        if comp3_score > 0:
            total_score += comp3_score
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # ----------------------------------------------------------------
    # Component 4: Slides 6-7 notes contain expected text (0.20 points)
    # Per slide: 0.10 points, each bullet within that slide is proportional.
    # This FAILS on initial_env (empty notes) and PASSES on golden_env.
    # ----------------------------------------------------------------
    try:
        comp4_score = 0.0
        for slide_num in range(6, 8):
            notes_text = all_notes[slide_num]
            expected_bullets = EXPECTED_NOTES[slide_num]
            n_found = count_matching_bullets(notes_text, expected_bullets)
            n_total = len(expected_bullets)
            per_bullet_pts = 0.10 / n_total
            if n_found == n_total:
                print(f"PASS: Component 4 — Slide {slide_num}: all {n_total} bullets found")
                comp4_score += 0.10
            elif n_found > 0:
                slide_pts = round(n_found * per_bullet_pts, 4)
                print(f"PARTIAL: Component 4 — Slide {slide_num}: {n_found}/{n_total} bullets ({slide_pts} pts)")
                comp4_score += slide_pts
            else:
                print(f"FAIL: Component 4 — Slide {slide_num}: 0/{n_total} bullets. Notes: {repr(notes_text[:60])}")
        comp4_score = round(min(comp4_score, 0.20), 4)
        print(f"Component 4 subtotal: {comp4_score}/0.20")
        if comp4_score > 0:
            total_score += comp4_score
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = round(min(total_score, 1.0), 4)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/Startup_Pitch.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
