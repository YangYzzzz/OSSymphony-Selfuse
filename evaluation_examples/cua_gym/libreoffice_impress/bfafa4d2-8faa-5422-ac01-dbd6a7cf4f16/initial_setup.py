"""
Initial Setup: Startup pitch deck with empty notes + talking points docx on Desktop
Task ID: osworld_multi_apps_impress_notes_import_005
Domain: libreoffice_impress

Creates:
  - /home/user/Startup_Pitch.pptx  — 7-slide startup pitch, all notes empty
  - /home/user/Desktop/pitch_talking_points.docx — 7 sections with bullet points
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from docx import Document
from docx.shared import Pt as DocxPt

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
TASK_PPTX = f'{WORKDIR}/Startup_Pitch.pptx'
TASK_DOCX = f'{DESKTOP}/pitch_talking_points.docx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_pptx():
    """Create a 7-slide startup pitch deck with empty notes."""
    prs = Presentation()

    # Slide data: (layout_idx, title, subtitle_or_content)
    slides_data = [
        {
            "layout": 0,
            "title": "NovaMind AI",
            "subtitle": "Redefining Intelligent Automation for Enterprise Teams",
        },
        {
            "layout": 1,
            "title": "The Problem",
            "bullets": [
                "Enterprise teams waste 40% of work hours on repetitive manual tasks",
                "Existing automation tools require specialized engineering knowledge",
                "High implementation costs and long deployment cycles deter SMBs",
                "Siloed systems prevent end-to-end workflow automation",
            ],
        },
        {
            "layout": 1,
            "title": "Our Solution",
            "bullets": [
                "NovaMind AI: a no-code platform that automates cross-app workflows",
                "Natural language task definition — describe what you want, we automate it",
                "Pre-built connectors for 150+ enterprise SaaS tools",
                "Average deployment time: 2 hours vs. industry average of 3 weeks",
            ],
        },
        {
            "layout": 1,
            "title": "Market Opportunity",
            "bullets": [
                "Total addressable market: $47B by 2027 (IDC, 2024)",
                "Serviceable addressable market: $12B — mid-market enterprise segment",
                "Current market penetration: <5% — massive greenfield opportunity",
                "Growth rate: 28% CAGR in intelligent process automation",
            ],
        },
        {
            "layout": 1,
            "title": "Traction & Metrics",
            "bullets": [
                "120 enterprise customers in 14 months since launch",
                "ARR: $3.2M — 3x growth quarter-over-quarter",
                "Net Revenue Retention: 138%",
                "Pilot-to-paid conversion rate: 74%",
                "Notable customers: Meridian Health, TerraLogix, Cascade Financial",
            ],
        },
        {
            "layout": 1,
            "title": "Business Model & Financials",
            "bullets": [
                "SaaS subscription: $2,500–$25,000/month based on seats and workflows",
                "Professional services: one-time onboarding fee ($5,000–$15,000)",
                "Gross margin: 78% at current scale",
                "Path to profitability: Q3 2026 at $8M ARR",
            ],
        },
        {
            "layout": 1,
            "title": "The Ask & Use of Funds",
            "bullets": [
                "Raising $12M Series A",
                "40% — Sales & marketing expansion (5 new markets)",
                "30% — R&D: Advanced AI reasoning and enterprise integrations",
                "20% — Customer success and onboarding scale",
                "10% — G&A and infrastructure",
            ],
        },
    ]

    for i, slide_info in enumerate(slides_data):
        layout_idx = slide_info["layout"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = slide_info["title"]

        # Set content
        if layout_idx == 0:
            # Title slide: set subtitle placeholder
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_info["subtitle"]
        else:
            # Content slide: set bullets
            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for j, bullet in enumerate(slide_info["bullets"]):
                    if j == 0:
                        para = tf.paragraphs[0]
                    else:
                        para = tf.add_paragraph()
                    para.text = bullet
                    para.level = 0

        # Notes are intentionally left EMPTY for initial state
        # (Do NOT access notes_slide to avoid creating a notes placeholder)

    prs.save(TASK_PPTX)
    print(f'Startup_Pitch.pptx created: {TASK_PPTX}')


def create_docx():
    """Create the talking points document with 7 sections of bullet points."""
    os.makedirs(DESKTOP, exist_ok=True)

    doc = Document()
    doc.add_heading('Pitch Talking Points', level=0)

    talking_points = [
        {
            "header": "Slide 1:",
            "bullets": [
                "Welcome the audience and introduce yourself as CEO and co-founder",
                "NovaMind AI was born out of personal frustration automating enterprise workflows",
                "We are here today because we believe AI-driven automation should be accessible to every business",
                "Our team has 40+ years of combined experience in enterprise software and AI",
            ],
        },
        {
            "header": "Slide 2:",
            "bullets": [
                "Every day, knowledge workers spend nearly half their time on tasks that could be automated",
                "Legacy RPA tools are brittle, expensive, and require dedicated engineering teams to maintain",
                "Small and mid-sized businesses are priced out of enterprise automation entirely",
                "The result: billions of dollars in lost productivity and missed revenue opportunities annually",
                "We validated this with 200+ interviews across finance, healthcare, and logistics sectors",
            ],
        },
        {
            "header": "Slide 3:",
            "bullets": [
                "NovaMind AI lets any business user describe a workflow in plain English and deploy it instantly",
                "Our AI engine maps intent to actions across connected apps with no coding required",
                "Security-first architecture: SOC 2 Type II certified, data never leaves your cloud environment",
                "We integrate with the tools your teams already use — no migration needed",
            ],
        },
        {
            "header": "Slide 4:",
            "bullets": [
                "The intelligent process automation market is growing faster than analysts predicted",
                "Regulation changes in finance and healthcare are accelerating digital transformation",
                "Our sweet spot: companies with 50–500 employees that have outgrown spreadsheets",
                "Competitive moat: proprietary workflow graph model trained on 10M+ enterprise tasks",
            ],
        },
        {
            "header": "Slide 5:",
            "bullets": [
                "We are not pre-revenue — we have paying customers proving real-world value",
                "Meridian Health reduced claims processing time by 67% in 90 days",
                "TerraLogix automated their entire vendor onboarding process, saving 1,200 hours per quarter",
                "Our NPS score is 72 — exceptional for enterprise software",
                "Pipeline: 45 active enterprise pilots expected to convert in Q2 2026",
            ],
        },
        {
            "header": "Slide 6:",
            "bullets": [
                "Simple, predictable pricing aligned with customer value delivery",
                "We expand naturally as customers automate more workflows — land and expand motion",
                "Customer acquisition cost: $8,400; lifetime value: $142,000 — LTV:CAC ratio of 17:1",
                "Unit economics strengthen at scale as AI model accuracy improves with more data",
            ],
        },
        {
            "header": "Slide 7:",
            "bullets": [
                "We are raising $12M to accelerate what is already working",
                "Series A will fund expansion into the UK, Germany, and Australia",
                "Product roadmap: AI Copilot feature launching Q4 2026 — projected to double ARPU",
                "We are in conversations with three strategic VCs and welcome your partnership",
                "Thank you — we are excited to answer your questions and explore how we grow together",
            ],
        },
    ]

    for section in talking_points:
        # Section header (e.g., "Slide 1:")
        heading_para = doc.add_heading(section["header"], level=1)

        # Bullet points
        for bullet_text in section["bullets"]:
            para = doc.add_paragraph(style='List Bullet')
            para.text = bullet_text

        # Add a blank line between sections
        doc.add_paragraph()

    doc.save(TASK_DOCX)
    print(f'pitch_talking_points.docx created: {TASK_DOCX}')


def main():
    create_pptx()
    create_docx()

    # GUI-ready startup: open the pptx in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{TASK_PPTX}"', delay_sec=3.0)
    # Also open the docx so the user can reference it easily
    launch_gui(f'libreoffice --writer "{TASK_DOCX}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress and Writer with DISPLAY=:0')


main()
