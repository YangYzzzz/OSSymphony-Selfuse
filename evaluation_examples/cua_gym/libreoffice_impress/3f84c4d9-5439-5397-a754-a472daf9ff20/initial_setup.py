"""
Initial Setup: Insert video on slide 6 of a product demo presentation
Task ID: impress_gf1_037
Domain: libreoffice_impress

Creates an 8-slide Product Demo presentation with slide 6 titled 'Live Demo'
and an empty content area. Also creates a dummy demo_clip.mp4 on the Desktop.
"""

import os
import shlex
import subprocess
import time
import struct
import shutil

WORKDIR = '/home/user'
TASK_ID = 'impress_gf1_037'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
VIDEO_PATH = f'{WORKDIR}/Desktop/demo_clip.mp4'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_dummy_mp4(path):
    """Create a minimal valid MP4 file (1920x1080 metadata) using ffmpeg if available, else raw bytes."""
    os.makedirs(os.path.dirname(path), exist_ok=True)
    # Try ffmpeg first for a proper MP4
    try:
        subprocess.run(
            [
                'ffmpeg', '-y',
                '-f', 'lavfi', '-i', 'color=c=blue:s=1920x1080:d=5:r=24',
                '-f', 'lavfi', '-i', 'anullsrc=r=44100:cl=stereo',
                '-t', '5',
                '-c:v', 'libx264', '-preset', 'ultrafast', '-crf', '51',
                '-c:a', 'aac', '-b:a', '32k',
                '-shortest',
                path
            ],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            timeout=30,
        )
        print(f'Created MP4 via ffmpeg: {path}')
        return
    except (FileNotFoundError, subprocess.CalledProcessError, subprocess.TimeoutExpired):
        pass

    # Fallback: write a minimal MP4 file with ftyp + moov boxes
    # This won't play but is recognized as MP4 by LibreOffice
    with open(path, 'wb') as f:
        # ftyp box
        ftyp = b'\x00\x00\x00\x18ftypmp42\x00\x00\x00\x00mp42isom'
        f.write(ftyp)
        # minimal moov box
        moov = b'\x00\x00\x00\x08moov'
        f.write(moov)
    print(f'Created minimal MP4 fallback: {path}')


def create_initial():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Product Demo"
    slide1.placeholders[1].text = "Q2 2025 Launch Overview\nPrepared by the Product Team"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Product overview and market positioning"
    items2 = [
        "Key features deep dive",
        "Technical architecture walkthrough",
        "Live product demonstration",
        "Customer success stories",
        "Q&A session",
    ]
    for item in items2:
        p = body2.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 3: Product Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Product Overview"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "NexaFlow Analytics Platform"
    features = [
        "Real-time data processing with sub-100ms latency",
        "AI-powered anomaly detection across 50+ signal types",
        "Enterprise-grade security with SOC2 Type II compliance",
        "Seamless integration with Salesforce, HubSpot, and Snowflake",
        "Custom dashboard builder with 30+ visualization templates",
    ]
    for feat in features:
        p = body3.add_paragraph()
        p.text = feat
        p.level = 1

    # --- Slide 4: Key Features ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Key Features"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Smart Pipeline Orchestration"
    kf_items = [
        "Drag-and-drop ETL pipeline builder",
        "Auto-scaling compute clusters (AWS, GCP, Azure)",
        "Version-controlled data transformations with Git integration",
        "Scheduled and event-driven pipeline triggers",
        "Built-in data quality checks and alerting",
    ]
    for kf in kf_items:
        p = body4.add_paragraph()
        p.text = kf
        p.level = 1

    # --- Slide 5: Technical Architecture ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Technical Architecture"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Microservices-based architecture deployed on Kubernetes"
    arch_items = [
        "API Gateway: Kong with rate limiting and OAuth2",
        "Message Queue: Apache Kafka for event streaming",
        "Storage: PostgreSQL + Redis + S3-compatible object store",
        "Compute: Spark clusters managed by Airflow DAGs",
        "Monitoring: Prometheus + Grafana with custom SLO dashboards",
    ]
    for ai in arch_items:
        p = body5.add_paragraph()
        p.text = ai
        p.level = 1

    # --- Slide 6: Live Demo (title only, empty content area) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide6.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Live Demo"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
    # No other content - empty area for video insertion

    # --- Slide 7: Customer Testimonials ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Customer Testimonials"
    body7 = slide7.placeholders[1].text_frame
    body7.text = '"NexaFlow reduced our reporting time by 73% in the first quarter."'
    testimonials = [
        "— Rachel Torres, VP of Data, Meridian Financial",
        "",
        '"The anomaly detection alone saved us $2.1M in prevented fraud last year."',
        "— David Kim, CTO, Pacific Commerce Group",
        "",
        '"Best-in-class onboarding. Our team was productive within two weeks."',
        "— Lisa Andersen, Head of Analytics, Nordic Health Systems",
    ]
    for t in testimonials:
        p = body7.add_paragraph()
        p.text = t
        p.level = 0

    # --- Slide 8: Thank You / Q&A ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    txBox8 = slide8.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(2))
    tf8 = txBox8.text_frame
    p8 = tf8.paragraphs[0]
    p8.text = "Thank You"
    p8.alignment = PP_ALIGN.CENTER
    run8 = p8.runs[0]
    run8.font.size = Pt(44)
    run8.font.bold = True
    run8.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    p8b = tf8.add_paragraph()
    p8b.text = "Questions & Answers"
    p8b.alignment = PP_ALIGN.CENTER
    run8b = p8b.runs[0]
    run8b.font.size = Pt(28)
    run8b.font.color.rgb = RGBColor(0x4A, 0x4A, 0x4A)

    p8c = tf8.add_paragraph()
    p8c.text = "Contact: product-team@nexaflow.io"
    p8c.alignment = PP_ALIGN.CENTER
    run8c = p8c.runs[0]
    run8c.font.size = Pt(16)
    run8c.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Create demo video file
    create_dummy_mp4(VIDEO_PATH)

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
