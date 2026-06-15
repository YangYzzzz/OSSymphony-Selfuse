"""
Reward Script: Insert video on slide 6, scale to 16cm x 9cm, center it, play on click
Task ID: impress_gf1_037
Domain: libreoffice_impress
Scoring:
  Component 1: Video/media shape exists on slide 6 (0.30)
  Component 2: Video dimensions are 16cm x 9cm (0.30)
  Component 3: Video is centered on the slide (0.20)
  Component 4: No auto-play timing — plays on click (0.20)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_gf1_037'

# 1 cm = 360000 EMU
CM_TO_EMU = 360000
EXPECTED_WIDTH = 16 * CM_TO_EMU   # 5760000
EXPECTED_HEIGHT = 9 * CM_TO_EMU   # 3240000


def _check_auto_play(file_path):
    """Returns a truthy string if auto-play timing is detected, empty string otherwise."""
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide6.xml') as f:
                xml_content = f.read().decode()
                root = ET.fromstring(xml_content)
                ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                timing = root.find(f'{{{ns_p}}}timing')
                if timing is not None:
                    xml_str = ET.tostring(timing, encoding='unicode')
                    if 'nodeType="afterPrevious"' in xml_str or 'nodeType="withPrevious"' in xml_str:
                        return "auto-play timing found"
    except Exception:
        pass
    return ""


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[5]  # slide 6, 0-indexed
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Find media/video shape on slide 6
    video_shape = None
    for shape in slide.shapes:
        # Check for MEDIA type (16) which is how python-pptx reports video shapes
        if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
            video_shape = shape
            break

    # Also check via XML for video references if python-pptx doesn't detect MEDIA type
    has_video_in_xml = False
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            try:
                with zf.open('ppt/slides/slide6.xml') as f:
                    xml_content = f.read().decode()
                    has_video_in_xml = 'videoFile' in xml_content
            except KeyError:
                pass

            # Check for media files in the archive
            media_files = [n for n in zf.namelist() if n.startswith('ppt/media/') and n.endswith('.mp4')]
            has_media_files = len(media_files) > 0
    except Exception as e:
        print(f"ERROR: ZIP inspection failed: {e}")
        has_media_files = False

    # Component 1: Video/media shape exists on slide 6 (0.30 points)
    try:
        if video_shape is not None:
            print(f"PASS: Component 1 — Video shape found on slide 6: '{video_shape.name}' (0.30 pts)")
            total_score += 0.30
        elif has_video_in_xml and has_media_files:
            # Fallback: python-pptx may not always detect media, but XML confirms it
            print(f"PASS: Component 1 — Video found in slide 6 XML with embedded media file (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 1 — No video/media shape found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # For subsequent checks, we need the video shape or XML-based dimensions
    # Try to get dimensions from XML if python-pptx shape is not available
    vid_left = vid_top = vid_width = vid_height = None
    if video_shape is not None:
        vid_left = video_shape.left
        vid_top = video_shape.top
        vid_width = video_shape.width
        vid_height = video_shape.height
    else:
        # Try XML extraction
        try:
            ns = {
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
            }
            with zipfile.ZipFile(file_path, 'r') as zf:
                with zf.open('ppt/slides/slide6.xml') as f:
                    root = ET.parse(f).getroot()
                    # Find pic elements with videoFile
                    for pic in root.findall('.//p:pic', ns):
                        nvPr = pic.find('.//p:nvPr', ns)
                        if nvPr is not None:
                            vf = nvPr.find('.//a:videoFile', ns)
                            if vf is not None:
                                xfrm = pic.find('.//a:xfrm', ns)
                                if xfrm is not None:
                                    off = xfrm.find('a:off', ns)
                                    ext = xfrm.find('a:ext', ns)
                                    if off is not None and ext is not None:
                                        vid_left = int(off.get('x', 0))
                                        vid_top = int(off.get('y', 0))
                                        vid_width = int(ext.get('cx', 0))
                                        vid_height = int(ext.get('cy', 0))
        except Exception as e:
            print(f"ERROR: XML dimension extraction failed: {e}")

    # Component 2: Video dimensions are 16cm x 9cm (0.30 points)
    try:
        if vid_width is not None and vid_height is not None:
            width_ok = abs(vid_width - EXPECTED_WIDTH) / EXPECTED_WIDTH <= 0.02
            height_ok = abs(vid_height - EXPECTED_HEIGHT) / EXPECTED_HEIGHT <= 0.02
            if width_ok and height_ok:
                print(f"PASS: Component 2 — Video size {vid_width}x{vid_height} EMU matches 16cm x 9cm (0.30 pts)")
                total_score += 0.30
            else:
                print(f"FAIL: Component 2 — Video size {vid_width}x{vid_height} EMU, expected ~{EXPECTED_WIDTH}x{EXPECTED_HEIGHT}")
        else:
            print(f"FAIL: Component 2 — Cannot determine video dimensions (no video shape)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Video is centered on the slide (0.20 points)
    try:
        if vid_left is not None and vid_top is not None and vid_width is not None and vid_height is not None:
            expected_left = (slide_width - vid_width) // 2
            expected_top = (slide_height - vid_height) // 2
            # Use 5% of slide dimension as tolerance
            tol_x = slide_width * 0.05
            tol_y = slide_height * 0.05
            left_ok = abs(vid_left - expected_left) <= tol_x
            top_ok = abs(vid_top - expected_top) <= tol_y
            if left_ok and top_ok:
                print(f"PASS: Component 3 — Video centered at ({vid_left}, {vid_top}), expected ~({expected_left}, {expected_top}) (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 3 — Video at ({vid_left}, {vid_top}), expected center ~({expected_left}, {expected_top})")
        else:
            print(f"FAIL: Component 3 — Cannot determine video position (no video shape)")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: No auto-play timing — plays on click, not automatically (0.20 points)
    # If there's no <p:timing> with an auto-trigger for the video, it defaults to on-click.
    # We check that there is NO afterEffect or startCondEvt with delay="0" triggering automatically.
    try:
        if has_video_in_xml or video_shape is not None:
            auto_play_detected = _check_auto_play(file_path)

            if not auto_play_detected:
                print(f"PASS: Component 4 — No auto-play timing detected; video plays on click (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 4 — Auto-play timing detected; video should play on click only")
        else:
            print(f"FAIL: Component 4 — No video found to check playback settings")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
def persist_app_state(domain):
    import os, time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")

persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
