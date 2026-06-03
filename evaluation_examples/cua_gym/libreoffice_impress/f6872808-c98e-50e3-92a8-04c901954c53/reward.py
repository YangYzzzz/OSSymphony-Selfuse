"""
Reward Script: Change Management Presentation (ERP Migration)
Task ID: impress_wf_083
Domain: libreoffice_impress
Scoring:
  Component 1:  File exists on Desktop (0.05)
  Component 2:  Exactly 12 slides (0.15)
  Component 3:  Slide 1 title text (0.10)
  Component 4:  Slide 4 ADKAR - 5 chevron shapes (0.10)
  Component 5:  Slide 5 impact assessment table (0.10)
  Component 6:  Slide 7 Gantt bar shapes (0.10)
  Component 7:  Slide 9 - 4 arc gauge shapes (0.10)
  Component 8:  Slide 10 risk register table with RAG (0.10)
  Component 9:  Slide 12 - 90-day action plan table (0.10)
  Component 10: Colors E65100 and 1565C0 used (0.10)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_083'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)

    # Component 1: File exists on Desktop (0.05 points)
    # This is task-introduced: initial_env has NO Change_Management.pptx on Desktop
    try:
        if os.path.exists(file_path):
            print(f"PASS: Component 1 - File exists at {file_path} (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 1 - File not found at {file_path}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Exactly 12 slides (0.15 points)
    try:
        num_slides = len(slides)
        if num_slides == 12:
            print(f"PASS: Component 2 - Exactly 12 slides found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 - Expected 12 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Slide 1 title contains "Change Management Plan" and "ERP Migration" (0.10 points)
    try:
        if len(slides) >= 1:
            slide1_text = ""
            for shape in slides[0].shapes:
                if shape.has_text_frame:
                    slide1_text += " " + shape.text_frame.text
            slide1_lower = slide1_text.lower()
            has_cmp = "change management plan" in slide1_lower
            has_erp = "erp migration" in slide1_lower
            if has_cmp and has_erp:
                print(f"PASS: Component 3 - Slide 1 has both title phrases (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 3 - Slide 1 text missing required phrases. "
                      f"CMP={has_cmp}, ERP={has_erp}. Text: {slide1_text[:100]}")
        else:
            print("FAIL: Component 3 - No slides in presentation")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Slide 4 has 5 chevron/arrow shapes for ADKAR model (0.10 points)
    # Task says 5 connected arrows: Awareness, Desire, Knowledge, Ability, Reinforcement
    try:
        if len(slides) >= 4:
            slide4 = slides[3]
            chevron_count = 0
            adkar_words = {"awareness", "desire", "knowledge", "ability", "reinforcement"}
            adkar_found = set()
            for shape in slide4.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    # Chevrons (type 52) are arrow-like shapes
                    try:
                        if shape.auto_shape_type is not None and shape.auto_shape_type == 52:  # CHEVRON
                            chevron_count += 1
                            if shape.has_text_frame:
                                txt = shape.text_frame.text.strip().lower()
                                for word in adkar_words:
                                    if word in txt:
                                        adkar_found.add(word)
                    except:
                        pass
            if chevron_count >= 5 and len(adkar_found) >= 5:
                print(f"PASS: Component 4 - Slide 4 has {chevron_count} chevrons with ADKAR labels (0.10 pts)")
                total_score += 0.10
            elif chevron_count >= 5:
                print(f"PARTIAL: Component 4 - {chevron_count} chevrons but only {len(adkar_found)}/5 ADKAR labels (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 4 - Expected 5 chevrons, found {chevron_count}. ADKAR labels: {adkar_found}")
        else:
            print("FAIL: Component 4 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Slide 5 has impact assessment table with correct headers (0.10 points)
    # Headers: Department, Impact Level, Readiness
    try:
        if len(slides) >= 5:
            slide5 = slides[4]
            table_found = False
            for shape in slide5.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    headers = [table.cell(0, c).text.strip().lower() for c in range(min(3, len(table.columns)))]
                    expected = ["department", "impact level", "readiness"]
                    if all(exp in headers for exp in expected):
                        if len(table.rows) >= 4:
                            print(f"PASS: Component 5 - Impact assessment table with correct headers, {len(table.rows)} rows (0.10 pts)")
                            total_score += 0.10
                        else:
                            print(f"PARTIAL: Component 5 - Correct headers but only {len(table.rows)} rows (0.05 pts)")
                            total_score += 0.05
                        table_found = True
                    else:
                        print(f"FAIL: Component 5 - Table headers mismatch. Found: {headers}")
                        table_found = True
            if not table_found:
                print("FAIL: Component 5 - No table found on slide 5")
        else:
            print("FAIL: Component 5 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    # Component 6: Slide 7 has horizontal bar shapes (Gantt chart) (0.10 points)
    # Expect rounded rectangle bars as Gantt bars
    try:
        if len(slides) >= 7:
            slide7 = slides[6]
            rounded_rect_count = 0
            has_title = False
            for shape in slide7.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    try:
                        if shape.auto_shape_type is not None and shape.auto_shape_type == 5:  # ROUNDED_RECTANGLE
                            rounded_rect_count += 1
                    except:
                        pass
                if shape.has_text_frame:
                    txt = shape.text_frame.text.lower()
                    if "training" in txt and "schedule" in txt:
                        has_title = True
            if rounded_rect_count >= 5 and has_title:
                print(f"PASS: Component 6 - Slide 7 Gantt chart with {rounded_rect_count} bars and title (0.10 pts)")
                total_score += 0.10
            elif rounded_rect_count >= 3:
                print(f"PARTIAL: Component 6 - {rounded_rect_count} bars found (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 6 - Expected >=5 rounded rect bars, found {rounded_rect_count}, title={has_title}")
        else:
            print("FAIL: Component 6 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 6 - {e}")

    # Component 7: Slide 9 has 4 arc shapes (KPI gauges) (0.10 points)
    try:
        if len(slides) >= 9:
            slide9 = slides[8]
            arc_count = 0
            for shape in slide9.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    try:
                        if shape.auto_shape_type is not None and shape.auto_shape_type == 25:  # ARC
                            arc_count += 1
                    except:
                        pass
            if arc_count >= 4:
                print(f"PASS: Component 7 - Slide 9 has {arc_count} arc gauge shapes (0.10 pts)")
                total_score += 0.10
            elif arc_count >= 2:
                print(f"PARTIAL: Component 7 - Expected 4 arcs, found {arc_count} (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 7 - Expected 4 arc gauges, found {arc_count}")
        else:
            print("FAIL: Component 7 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 7 - {e}")

    # Component 8: Slide 10 has risk register table with RAG Status column (0.10 points)
    try:
        if len(slides) >= 10:
            slide10 = slides[9]
            table_found = False
            for shape in slide10.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    headers = [table.cell(0, c).text.strip().lower() for c in range(len(table.columns))]
                    has_rag = any("rag" in h or "status" in h for h in headers)
                    has_risk = any("risk" in h for h in headers)
                    if has_rag and has_risk and len(table.rows) >= 3:
                        print(f"PASS: Component 8 - Risk register table with RAG, {len(table.rows)} rows (0.10 pts)")
                        total_score += 0.10
                    elif has_rag or has_risk:
                        print(f"PARTIAL: Component 8 - Table found but incomplete. RAG={has_rag}, Risk={has_risk} (0.05 pts)")
                        total_score += 0.05
                    else:
                        print(f"FAIL: Component 8 - Table headers lack RAG/Risk. Found: {headers}")
                    table_found = True
            if not table_found:
                print("FAIL: Component 8 - No table found on slide 10")
        else:
            print("FAIL: Component 8 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 8 - {e}")

    # Component 9: Slide 12 has 90-day action plan table (0.10 points)
    try:
        if len(slides) >= 12:
            slide12 = slides[11]
            table_found = False
            has_title = False
            for shape in slide12.shapes:
                if shape.has_text_frame:
                    txt = shape.text_frame.text.lower()
                    if "90" in txt and "action" in txt:
                        has_title = True
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    if len(table.rows) >= 3 and len(table.columns) >= 3:
                        print(f"PASS: Component 9 - 90-day action plan table {len(table.rows)}x{len(table.columns)} (0.10 pts)")
                        total_score += 0.10
                        table_found = True
                    else:
                        print(f"FAIL: Component 9 - Table too small: {len(table.rows)}x{len(table.columns)}")
                        table_found = True
            if not table_found:
                print("FAIL: Component 9 - No table found on slide 12")
        else:
            print("FAIL: Component 9 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 9 - {e}")

    # Component 10: Colors E65100 (change-orange) and 1565C0 (blue) used (0.10 points)
    try:
        colors_found = set()
        for slide in slides:
            for shape in slide.shapes:
                # Check shape fill
                try:
                    if shape.fill and shape.fill.type == 1:  # solid fill
                        c = str(shape.fill.fore_color.rgb)
                        colors_found.add(c)
                except:
                    pass
                # Check text font colors
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                if run.font.color and run.font.color.type is not None:
                                    c = str(run.font.color.rgb)
                                    colors_found.add(c)
                            except:
                                pass

        has_orange = "E65100" in colors_found
        has_blue = "1565C0" in colors_found
        if has_orange and has_blue:
            print(f"PASS: Component 10 - Both theme colors found: E65100 and 1565C0 (0.10 pts)")
            total_score += 0.10
        elif has_orange or has_blue:
            print(f"PARTIAL: Component 10 - Only one color found. Orange={has_orange}, Blue={has_blue} (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 10 - Neither E65100 nor 1565C0 found. Colors: {colors_found}")
    except Exception as e:
        print(f"ERROR: Component 10 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook: save any unsaved GUI state
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/Desktop/Change_Management.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
