"""
Reward Script: Reposition title text frames on slides 2 and 3 to the bottom of each slide.
Task ID: osworld_impress_title_position_bottom_005
Domain: libreoffice_impress
Scoring:
  Component 1: Title text box on slide 2 is at the bottom of the slide (0.5 pts)
  Component 2: Title text box on slide 3 is at the bottom of the slide (0.5 pts)
  Total: 1.0
"""

import os

from pptx import Presentation
from pptx.util import Emu

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_position_bottom_005'

# The title text boxes on slides 2 and 3 are TEXT_BOX shapes (shape_type == 17),
# named 'Title 2' and 'Title 3' respectively.
# In the initial state, both have top=365760 (near the top of the slide).
# In the golden state, both have top=5943600 (near the bottom of the slide).
# The slide height is 6858000 EMU.
# "Bottom" is defined as: top > slide_height / 2 (i.e., top > 3429000).
# We use a strict threshold: the title box must be in the lower half of the slide.


def is_at_bottom(shape_top, shape_height, slide_height, threshold_ratio=0.5):
    """
    Returns True if the shape top is in the lower half of the slide.
    threshold_ratio=0.5 means top must be beyond the midpoint.
    """
    midpoint = slide_height * threshold_ratio
    return shape_top > midpoint


def find_title_textbox_on_slide(slide, slide_idx):
    """
    Find the title text box (TEXT_BOX with title text) on a given slide.
    Slide 2 (idx=1): name='Title 2', text='Financial Performance'
    Slide 3 (idx=2): name='Title 3', text='Market Expansion Strategy'
    Returns the shape or None.
    """
    # TEXT_BOX shape_type == 17
    TEXT_BOX_TYPE = 17
    for shape in slide.shapes:
        if shape.shape_type == TEXT_BOX_TYPE:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:  # non-empty title textbox
                    return shape
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slide_height = prs.slide_height
    num_slides = len(prs.slides)

    print(f"INFO: Slide height = {slide_height} EMU ({slide_height/914400:.2f} inches)")
    print(f"INFO: Number of slides = {num_slides}")

    if num_slides < 3:
        print(f"FAIL: Expected at least 3 slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Title text box on slide 2 is at the bottom (0.5 points)
    try:
        slide2 = prs.slides[1]  # 0-indexed, slide 2
        title_shape_s2 = find_title_textbox_on_slide(slide2, 1)

        if title_shape_s2 is None:
            print("FAIL: Component 1 — No title text box found on slide 2")
        else:
            top = title_shape_s2.top
            height = title_shape_s2.height
            text = title_shape_s2.text_frame.text.strip()
            print(f"INFO: Slide 2 title text box '{text}': top={top}, height={height}")

            if is_at_bottom(top, height, slide_height, threshold_ratio=0.5):
                print(f"PASS: Component 1 — Slide 2 title text box is at bottom (top={top} > midpoint={slide_height//2}) (0.5 pts)")
                total_score += 0.5
            else:
                print(f"FAIL: Component 1 — Slide 2 title text box NOT at bottom. top={top}, slide_height={slide_height}, midpoint={slide_height//2}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Title text box on slide 3 is at the bottom (0.5 points)
    try:
        slide3 = prs.slides[2]  # 0-indexed, slide 3
        title_shape_s3 = find_title_textbox_on_slide(slide3, 2)

        if title_shape_s3 is None:
            print("FAIL: Component 2 — No title text box found on slide 3")
        else:
            top = title_shape_s3.top
            height = title_shape_s3.height
            text = title_shape_s3.text_frame.text.strip()
            print(f"INFO: Slide 3 title text box '{text}': top={top}, height={height}")

            if is_at_bottom(top, height, slide_height, threshold_ratio=0.5):
                print(f"PASS: Component 2 — Slide 3 title text box is at bottom (top={top} > midpoint={slide_height//2}) (0.5 pts)")
                total_score += 0.5
            else:
                print(f"FAIL: Component 2 — Slide 3 title text box NOT at bottom. top={top}, slide_height={slide_height}, midpoint={slide_height//2}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in a given env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
