"""
Initial Setup: 5-slide minimalist presentation with title text frames at top of slides 2 and 3.
Task ID: osworld_impress_title_position_bottom_005
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_position_bottom_005'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # Standard widescreen 16:9 dimensions
    slide_width = prs.slide_width    # 9144000 EMU (10 inches)
    slide_height = prs.slide_height  # 5143500 EMU (5.63 inches)

    # ---- Slide 1: Title slide ----
    layout_title = prs.slide_layouts[0]  # Title Slide layout
    slide1 = prs.slides.add_slide(layout_title)
    slide1.shapes.title.text = "Annual Report 2024"
    # Set subtitle
    for ph in slide1.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = "Strategic Overview & Key Highlights"

    # Style slide 1 background
    fill1 = slide1.background.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

    # ---- Slide 2: Content slide with title at TOP ----
    layout_blank = prs.slide_layouts[5]  # Blank layout
    slide2 = prs.slides.add_slide(layout_blank)

    # Background
    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Title text box at TOP of slide (top ~0.4 inches from top)
    title2_left = Inches(0.8)
    title2_top = Inches(0.4)
    title2_width = Inches(8.4)
    title2_height = Inches(0.8)
    txb2 = slide2.shapes.add_textbox(title2_left, title2_top, title2_width, title2_height)
    txb2.name = "Title 2"
    tf2 = txb2.text_frame
    tf2.word_wrap = False
    p2 = tf2.paragraphs[0]
    p2.text = "Financial Performance"
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.runs[0]
    run2.font.name = "Calibri"
    run2.font.size = Pt(28)
    run2.font.bold = True
    run2.font.color.rgb = RGBColor(0x2E, 0x2E, 0x2E)

    # Content area: text box with bullet items
    content2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(3.5))
    cf2 = content2.text_frame
    cf2.word_wrap = True
    lines2 = [
        ("Total Revenue: $4.2 Billion (+12% YoY)", 0),
        ("Operating Margin: 18.3%", 0),
        ("Net Income: $768 Million", 0),
        ("Earnings Per Share: $3.47", 0),
        ("Free Cash Flow: $1.1 Billion", 0),
    ]
    for i, (line, lvl) in enumerate(lines2):
        if i == 0:
            para = cf2.paragraphs[0]
        else:
            para = cf2.add_paragraph()
        para.text = line
        para.level = lvl
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # ---- Slide 3: Content slide with title at TOP ----
    slide3 = prs.slides.add_slide(layout_blank)

    fill3 = slide3.background.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Title text box at TOP of slide
    title3_left = Inches(0.8)
    title3_top = Inches(0.4)
    title3_width = Inches(8.4)
    title3_height = Inches(0.8)
    txb3 = slide3.shapes.add_textbox(title3_left, title3_top, title3_width, title3_height)
    txb3.name = "Title 3"
    tf3 = txb3.text_frame
    tf3.word_wrap = False
    p3 = tf3.paragraphs[0]
    p3.text = "Market Expansion Strategy"
    p3.alignment = PP_ALIGN.LEFT
    run3 = p3.runs[0]
    run3.font.name = "Calibri"
    run3.font.size = Pt(28)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0x2E, 0x2E, 0x2E)

    # Content area
    content3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(3.5))
    cf3 = content3.text_frame
    cf3.word_wrap = True
    lines3 = [
        ("Entered 3 new international markets in Q3", 0),
        ("Asia-Pacific region: 24% revenue growth", 0),
        ("European partnerships signed: 7 enterprise deals", 0),
        ("North America market share increased to 31%", 0),
        ("Target: 15 new markets by end of 2025", 0),
    ]
    for i, (line, lvl) in enumerate(lines3):
        if i == 0:
            para = cf3.paragraphs[0]
        else:
            para = cf3.add_paragraph()
        para.text = line
        para.level = lvl
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # ---- Slide 4: Product overview ----
    slide4 = prs.slides.add_slide(layout_blank)

    fill4 = slide4.background.fill
    fill4.solid()
    fill4.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    txb4_title = slide4.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.8))
    tf4t = txb4_title.text_frame
    p4t = tf4t.paragraphs[0]
    p4t.text = "Product Portfolio Update"
    run4t = p4t.runs[0]
    run4t.font.name = "Calibri"
    run4t.font.size = Pt(28)
    run4t.font.bold = True
    run4t.font.color.rgb = RGBColor(0x2E, 0x2E, 0x2E)

    content4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(3.5))
    cf4 = content4.text_frame
    cf4.word_wrap = True
    lines4 = [
        ("CloudSync Pro: 2.1M active enterprise users", 0),
        ("DataInsight Analytics: Launched in 42 countries", 0),
        ("SecureVault 3.0: 99.99% uptime SLA achieved", 0),
        ("MobileConnect Suite: 8.7M downloads in 2024", 0),
        ("New: AI-Assist Module — Beta testing Q4 2024", 0),
    ]
    for i, (line, lvl) in enumerate(lines4):
        if i == 0:
            para = cf4.paragraphs[0]
        else:
            para = cf4.add_paragraph()
        para.text = line
        para.level = lvl
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # ---- Slide 5: Closing / Outlook ----
    slide5 = prs.slides.add_slide(layout_blank)

    fill5 = slide5.background.fill
    fill5.solid()
    fill5.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

    txb5_title = slide5.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.8))
    tf5t = txb5_title.text_frame
    p5t = tf5t.paragraphs[0]
    p5t.text = "2025 Outlook & Closing Remarks"
    run5t = p5t.runs[0]
    run5t.font.name = "Calibri"
    run5t.font.size = Pt(28)
    run5t.font.bold = True
    run5t.font.color.rgb = RGBColor(0x2E, 0x2E, 0x2E)

    content5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(3.0))
    cf5 = content5.text_frame
    cf5.word_wrap = True
    closing_lines = [
        ("Revenue guidance: $4.7B – $4.9B", 0),
        ("Planned headcount growth: +1,200 employees", 0),
        ("R&D investment: $620M (increase of 15%)", 0),
        ("Dividend increase: $1.20 per share quarterly", 0),
    ]
    for i, (line, lvl) in enumerate(closing_lines):
        if i == 0:
            para = cf5.paragraphs[0]
        else:
            para = cf5.add_paragraph()
        para.text = line
        para.level = lvl
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
