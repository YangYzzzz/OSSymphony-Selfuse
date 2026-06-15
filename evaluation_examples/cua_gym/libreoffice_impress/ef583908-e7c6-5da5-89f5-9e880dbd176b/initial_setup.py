"""
Initial Setup: Set slide 2 body text to right-aligned and slide 3 body text to center-aligned.
Task ID: osworld_impress_per_slide_alignment_004
Domain: libreoffice_impress

Creates a 6-slide marketing deck with all body text left-aligned.
Slides 2 and 3 have left-aligned body text (the agent must change them).
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_per_slide_alignment_004'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_paragraph_alignment(para, alignment):
    """Set paragraph alignment."""
    para.alignment = alignment


def create_initial():
    prs = Presentation()
    # Use standard 16:9 widescreen
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_layouts = prs.slide_layouts
    # Layout 0 = Title Slide, Layout 1 = Title + Content, Layout 5 = Blank

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(slide_layouts[0])
    slide1.shapes.title.text = "GlobalReach Marketing Solutions"
    subtitle = slide1.placeholders[1]
    subtitle.text = "Annual Strategy Presentation 2025"
    for para in subtitle.text_frame.paragraphs:
        set_paragraph_alignment(para, PP_ALIGN.LEFT)

    # --- Slide 2: Market Overview (body text left-aligned — MUST remain left) ---
    slide2 = prs.slides.add_slide(slide_layouts[1])
    slide2.shapes.title.text = "Market Overview"
    body2 = slide2.placeholders[1]
    tf2 = body2.text_frame
    tf2.clear()
    lines2 = [
        "Global digital advertising spend reached $621 billion in 2024",
        "Mobile-first campaigns now represent 68% of all impressions",
        "AI-driven targeting has improved conversion rates by 34%",
        "Social commerce is the fastest-growing channel at 29% YoY",
        "Personalization at scale remains the top priority for brands",
    ]
    for i, line in enumerate(lines2):
        if i == 0:
            para = tf2.paragraphs[0]
        else:
            para = tf2.add_paragraph()
        para.text = line
        para.alignment = PP_ALIGN.LEFT  # MUST be left — task requires agent to change to RIGHT
        for run in para.runs:
            run.font.size = Pt(18)

    # --- Slide 3: Campaign Performance (body text left-aligned — MUST remain left) ---
    slide3 = prs.slides.add_slide(slide_layouts[1])
    slide3.shapes.title.text = "Campaign Performance"
    body3 = slide3.placeholders[1]
    tf3 = body3.text_frame
    tf3.clear()
    lines3 = [
        "Q1 total impressions: 2.4 billion across all channels",
        "Click-through rate improved from 1.8% to 2.6% vs. prior year",
        "Customer acquisition cost reduced by 18% through optimization",
        "Brand awareness lifted 12 points in key demographics",
        "Return on ad spend (ROAS) achieved 4.2x industry benchmark",
    ]
    for i, line in enumerate(lines3):
        if i == 0:
            para = tf3.paragraphs[0]
        else:
            para = tf3.add_paragraph()
        para.text = line
        para.alignment = PP_ALIGN.LEFT  # MUST be left — task requires agent to change to CENTER
        for run in para.runs:
            run.font.size = Pt(18)

    # --- Slide 4: Target Audience Segments ---
    slide4 = prs.slides.add_slide(slide_layouts[1])
    slide4.shapes.title.text = "Target Audience Segments"
    body4 = slide4.placeholders[1]
    tf4 = body4.text_frame
    tf4.clear()
    segments = [
        "Millennials (25–40): Digital natives, value authenticity and sustainability",
        "Gen Z (18–24): Short-form video, peer influence, mobile-dominant",
        "Gen X (41–56): High purchasing power, loyalty to trusted brands",
        "Boomers (57–75): Growing online presence, health and lifestyle focus",
    ]
    for i, line in enumerate(segments):
        if i == 0:
            para = tf4.paragraphs[0]
        else:
            para = tf4.add_paragraph()
        para.text = line
        para.alignment = PP_ALIGN.LEFT
        for run in para.runs:
            run.font.size = Pt(18)

    # --- Slide 5: Strategic Initiatives ---
    slide5 = prs.slides.add_slide(slide_layouts[1])
    slide5.shapes.title.text = "Strategic Initiatives for 2025"
    body5 = slide5.placeholders[1]
    tf5 = body5.text_frame
    tf5.clear()
    initiatives = [
        "Expand programmatic buying to connected TV and streaming platforms",
        "Launch first-party data strategy ahead of cookie deprecation",
        "Invest in creator partnerships and micro-influencer networks",
        "Deploy dynamic creative optimization across all display formats",
        "Build cross-channel measurement framework with unified attribution",
    ]
    for i, line in enumerate(initiatives):
        if i == 0:
            para = tf5.paragraphs[0]
        else:
            para = tf5.add_paragraph()
        para.text = line
        para.alignment = PP_ALIGN.LEFT
        for run in para.runs:
            run.font.size = Pt(18)

    # --- Slide 6: Next Steps & Conclusion ---
    slide6 = prs.slides.add_slide(slide_layouts[1])
    slide6.shapes.title.text = "Next Steps & Conclusion"
    body6 = slide6.placeholders[1]
    tf6 = body6.text_frame
    tf6.clear()
    next_steps = [
        "Finalize Q2 budget allocation by March 31, 2025",
        "Schedule kickoff meetings with all channel leads",
        "Deliver updated creative briefs to agency partners",
        "Review and approve 2025 measurement plan",
        "Present progress update to executive team in April",
    ]
    for i, line in enumerate(next_steps):
        if i == 0:
            para = tf6.paragraphs[0]
        else:
            para = tf6.add_paragraph()
        para.text = line
        para.alignment = PP_ALIGN.LEFT
        for run in para.runs:
            run.font.size = Pt(18)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup — open the presentation in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
