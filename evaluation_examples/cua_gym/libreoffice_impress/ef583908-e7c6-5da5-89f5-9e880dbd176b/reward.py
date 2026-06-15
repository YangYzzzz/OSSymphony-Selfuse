"""
Reward Script: Set slide 2 body text to right-aligned and slide 3 body text to center-aligned.
Task ID: osworld_impress_per_slide_alignment_004
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): All paragraphs in slide 2 body (Content Placeholder 2) are RIGHT-aligned
  Component 2 (0.5): All paragraphs in slide 3 body (Content Placeholder 2) are CENTER-aligned
"""

import os
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_per_slide_alignment_004'


def get_body_content_shape(slide):
    """Return the first content/body placeholder that is NOT the title."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name != 'Title 1':
            return shape
    return None


def check_all_paragraphs_alignment(shape, expected_alignment):
    """
    Check that all non-empty paragraphs in the shape's text frame have the expected alignment.
    Returns (all_pass, total_paragraphs_checked, passing_count)
    """
    paras = shape.text_frame.paragraphs
    non_empty = [p for p in paras if p.text.strip()]
    if not non_empty:
        return False, 0, 0
    passing = sum(1 for p in non_empty if p.alignment == expected_alignment)
    return passing == len(non_empty), len(non_empty), passing


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Check slide count is 6
    num_slides = len(prs.slides)
    if num_slides < 3:
        print(f"CRITICAL: Expected at least 3 slides, found {num_slides}. Aborting.")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide 2 body text paragraphs are RIGHT-aligned (0.5 points)
    # This FAILS on initial (LEFT alignment) and PASSES on golden (RIGHT alignment)
    try:
        slide2 = prs.slides[1]  # 0-indexed
        body_shape2 = get_body_content_shape(slide2)
        if body_shape2 is None:
            print("FAIL: Component 1 — No body content shape found on slide 2")
        else:
            all_pass, total, passing = check_all_paragraphs_alignment(body_shape2, PP_ALIGN.RIGHT)
            if all_pass and total > 0:
                print(f"PASS: Component 1 — Slide 2 body text is RIGHT-aligned "
                      f"({passing}/{total} paragraphs, 0.5 pts)")
                total_score += 0.5
            else:
                # Show actual alignments for debugging
                actual_alignments = [p.alignment for p in body_shape2.text_frame.paragraphs if p.text.strip()]
                print(f"FAIL: Component 1 — Slide 2 body text NOT fully right-aligned. "
                      f"{passing}/{total} paragraphs pass. "
                      f"Actual alignments: {actual_alignments}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 3 body text paragraphs are CENTER-aligned (0.5 points)
    # This FAILS on initial (LEFT alignment) and PASSES on golden (CENTER alignment)
    try:
        slide3 = prs.slides[2]  # 0-indexed
        body_shape3 = get_body_content_shape(slide3)
        if body_shape3 is None:
            print("FAIL: Component 2 — No body content shape found on slide 3")
        else:
            all_pass, total, passing = check_all_paragraphs_alignment(body_shape3, PP_ALIGN.CENTER)
            if all_pass and total > 0:
                print(f"PASS: Component 2 — Slide 3 body text is CENTER-aligned "
                      f"({passing}/{total} paragraphs, 0.5 pts)")
                total_score += 0.5
            else:
                actual_alignments = [p.alignment for p in body_shape3.text_frame.paragraphs if p.text.strip()]
                print(f"FAIL: Component 2 — Slide 3 body text NOT fully center-aligned. "
                      f"{passing}/{total} paragraphs pass. "
                      f"Actual alignments: {actual_alignments}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the given env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
