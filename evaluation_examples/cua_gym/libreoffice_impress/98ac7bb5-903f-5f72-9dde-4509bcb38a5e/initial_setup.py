"""
Initial Setup: Create a checklist presentation with default round bullets
Task ID: impstruct_034
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impstruct_034'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_bullet_paragraph(tf, text, level=0, font_size=Pt(18), is_first=False):
    """Add a paragraph with a default round bullet."""
    if is_first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.level = level
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Set default round bullet explicitly
    pPr = p._p.get_or_add_pPr()
    # Remove any existing bullet elements
    for tag in [qn('a:buChar'), qn('a:buNone'), qn('a:buAutoNum')]:
        for existing in pPr.findall(tag):
            pPr.remove(existing)
    buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})  # round bullet
    pPr.append(buChar)
    # Set bullet size to match text
    buSzPct = pPr.makeelement(qn('a:buSzPct'), {'val': '100000'})  # 100%
    pPr.append(buSzPct)

    return p


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title
    txBox = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Q2 2025 Product Launch Checklist"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6D)

    # Subtitle
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Marketing & Operations Team"
    run2.font.size = Pt(20)
    run2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "Last Updated: March 28, 2025"
    run3.font.size = Pt(14)
    run3.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    # ---- Slide 2: Pre-Launch Tasks ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title
    txTitle2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(1))
    tf_t2 = txTitle2.text_frame
    r = tf_t2.paragraphs[0].add_run()
    r.text = "Pre-Launch Tasks"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6D)

    # Bullet list
    txBullets2 = slide2.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(10), Inches(5))
    tf_b2 = txBullets2.text_frame
    tf_b2.word_wrap = True

    bullets2 = [
        "Finalize product feature set and specifications",
        "Complete user acceptance testing for all modules",
        "Prepare marketing collateral and press releases",
        "Set up analytics dashboards for launch metrics",
        "Coordinate with design team on packaging mockups",
        "Review legal compliance documentation",
        "Schedule stakeholder review meeting for April 3",
    ]
    for i, text in enumerate(bullets2):
        add_bullet_paragraph(tf_b2, text, level=0, font_size=Pt(18), is_first=(i == 0))

    # ---- Slide 3: Launch Day Operations ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txTitle3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(1))
    tf_t3 = txTitle3.text_frame
    r3 = tf_t3.paragraphs[0].add_run()
    r3.text = "Launch Day Operations"
    r3.font.size = Pt(28)
    r3.font.bold = True
    r3.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6D)

    txBullets3 = slide3.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(10), Inches(5))
    tf_b3 = txBullets3.text_frame
    tf_b3.word_wrap = True

    bullets3 = [
        "Deploy production build to all regional servers by 6:00 AM",
        "Activate email campaign sequence for 150,000 subscribers",
        "Publish blog post and social media announcements",
        "Monitor server performance and error rate dashboards",
        "Brief customer support team on known issues and FAQ",
        "Send launch notification to enterprise clients via Slack",
    ]
    for i, text in enumerate(bullets3):
        add_bullet_paragraph(tf_b3, text, level=0, font_size=Pt(18), is_first=(i == 0))

    # ---- Slide 4: Post-Launch Follow-Up ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txTitle4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(1))
    tf_t4 = txTitle4.text_frame
    r4 = tf_t4.paragraphs[0].add_run()
    r4.text = "Post-Launch Follow-Up"
    r4.font.size = Pt(28)
    r4.font.bold = True
    r4.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6D)

    txBullets4 = slide4.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(10), Inches(5))
    tf_b4 = txBullets4.text_frame
    tf_b4.word_wrap = True

    bullets4 = [
        "Collect user feedback from first 48 hours of usage",
        "Analyze conversion funnel data and identify drop-off points",
        "Schedule retrospective meeting with engineering leads",
        "Prepare week-one performance report for leadership",
        "Prioritize bug fixes based on severity and user impact",
        "Update product roadmap with post-launch learnings",
        "Plan phase-two feature rollout timeline",
    ]
    for i, text in enumerate(bullets4):
        add_bullet_paragraph(tf_b4, text, level=0, font_size=Pt(18), is_first=(i == 0))

    # ---- Slide 5: Key Contacts ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txTitle5 = slide5.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(1))
    tf_t5 = txTitle5.text_frame
    r5 = tf_t5.paragraphs[0].add_run()
    r5.text = "Key Contacts"
    r5.font.size = Pt(28)
    r5.font.bold = True
    r5.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6D)

    txInfo5 = slide5.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(10), Inches(5))
    tf_i5 = txInfo5.text_frame
    tf_i5.word_wrap = True

    contacts = [
        ("Project Lead:", "Sarah Chen — sarah.chen@company.com"),
        ("Engineering:", "Marcus Johnson — marcus.j@company.com"),
        ("Marketing:", "Priya Sharma — priya.s@company.com"),
        ("Customer Support:", "David Kim — david.kim@company.com"),
    ]
    for i, (role, detail) in enumerate(contacts):
        if i == 0:
            p = tf_i5.paragraphs[0]
        else:
            p = tf_i5.add_paragraph()
        run_role = p.add_run()
        run_role.text = role + " "
        run_role.font.size = Pt(18)
        run_role.font.bold = True
        run_role.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6D)
        run_detail = p.add_run()
        run_detail.text = detail
        run_detail.font.size = Pt(18)
        run_detail.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
