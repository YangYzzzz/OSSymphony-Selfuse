"""
Reward Script: Apply underline and dark red (#8B0000) color to all text on slide 4
Task ID: osworld_impress_underline_darkred_table_003
Domain: libreoffice_impress
Scoring:
  Component 1: Title text underline applied (0.2 pts)
  Component 2: Title text color set to #8B0000 (0.2 pts)
  Component 3: All body text runs have underline applied (0.3 pts)
  Component 4: All body text runs have color set to #8B0000 (0.3 pts)
  Total: 1.0
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_underline_darkred_table_003'

TARGET_COLOR = RGBColor(0x8B, 0x00, 0x00)  # dark red #8B0000
TARGET_COLOR_STR = '8B0000'


def verify_task(file_path):
    """
    Verify that all text on slide 4 has underline formatting and dark red (#8B0000) color.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Ensure the presentation has at least 4 slides
    if len(prs.slides) < 4:
        print(f"CRITICAL: Expected at least 4 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # slide 4 (0-indexed)

    # Collect shapes: title (shape 0) and content (shape 1)
    title_shape = None
    content_shape = None

    for shape in slide.shapes:
        if shape.has_text_frame:
            if shape.name.startswith('Title') or (hasattr(shape, 'placeholder_format')
                    and shape.placeholder_format is not None
                    and shape.placeholder_format.idx == 0):
                title_shape = shape
            else:
                content_shape = shape

    # Component 1: Title text underline applied (0.2 pts)
    # The title on slide 4 initially has underline=None (no underline). Must become True.
    try:
        if title_shape is None:
            print("FAIL: Component 1 — Title shape not found on slide 4")
        else:
            title_runs = [
                run
                for para in title_shape.text_frame.paragraphs
                for run in para.runs
                if (run.text or "").strip()
            ]
            if not title_runs:
                print("FAIL: Component 1 — No runs in title shape")
            else:
                all_underlined = all(run.font.underline is True for run in title_runs)
                if all_underlined:
                    print(f"PASS: Component 1 — Title underline=True for all {len(title_runs)} run(s) (0.2 pts)")
                    total_score += 0.2
                else:
                    underline_vals = [run.font.underline for run in title_runs]
                    print(f"FAIL: Component 1 — Title underline values: {underline_vals}, expected all True")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Title text color set to #8B0000 (0.2 pts)
    # The title on slide 4 initially has color 1F397C. Must become 8B0000.
    try:
        if title_shape is None:
            print("FAIL: Component 2 — Title shape not found on slide 4")
        else:
            title_runs = [
                run
                for para in title_shape.text_frame.paragraphs
                for run in para.runs
                if (run.text or "").strip()
            ]
            if not title_runs:
                print("FAIL: Component 2 — No runs in title shape")
            else:
                color_ok = []
                for run in title_runs:
                    try:
                        if run.font.color.type is not None:
                            actual_color = str(run.font.color.rgb).upper()
                            color_ok.append(actual_color == TARGET_COLOR_STR)
                        else:
                            color_ok.append(False)
                    except Exception:
                        color_ok.append(False)

                if all(color_ok):
                    print(f"PASS: Component 2 — Title color=#8B0000 for all {len(title_runs)} run(s) (0.2 pts)")
                    total_score += 0.2
                else:
                    try:
                        actual_colors = [str(run.font.color.rgb) if run.font.color.type is not None else 'None' for run in title_runs]
                    except Exception:
                        actual_colors = ['Error']
                    print(f"FAIL: Component 2 — Title color values: {actual_colors}, expected all {TARGET_COLOR_STR}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: All body text runs have underline applied (0.3 pts)
    # The body on slide 4 initially has underline=False. Must become True.
    try:
        if content_shape is None:
            print("FAIL: Component 3 — Content shape not found on slide 4")
        else:
            body_runs = [
                run
                for para in content_shape.text_frame.paragraphs
                for run in para.runs
                if (run.text or "").strip()
            ]
            if not body_runs:
                print("FAIL: Component 3 — No runs in content/body shape")
            else:
                all_underlined = all(run.font.underline is True for run in body_runs)
                if all_underlined:
                    print(f"PASS: Component 3 — Body underline=True for all {len(body_runs)} run(s) (0.3 pts)")
                    total_score += 0.3
                else:
                    underline_vals = [run.font.underline for run in body_runs]
                    print(f"FAIL: Component 3 — Body underline values: {underline_vals}, expected all True")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: All body text runs have color set to #8B0000 (0.3 pts)
    # The body on slide 4 initially has color 333333. Must become 8B0000.
    try:
        if content_shape is None:
            print("FAIL: Component 4 — Content shape not found on slide 4")
        else:
            body_runs = [
                run
                for para in content_shape.text_frame.paragraphs
                for run in para.runs
                if (run.text or "").strip()
            ]
            if not body_runs:
                print("FAIL: Component 4 — No runs in content/body shape")
            else:
                color_ok = []
                for run in body_runs:
                    try:
                        if run.font.color.type is not None:
                            actual_color = str(run.font.color.rgb).upper()
                            color_ok.append(actual_color == TARGET_COLOR_STR)
                        else:
                            color_ok.append(False)
                    except Exception:
                        color_ok.append(False)

                if all(color_ok):
                    print(f"PASS: Component 4 — Body color=#8B0000 for all {len(body_runs)} run(s) (0.3 pts)")
                    total_score += 0.3
                else:
                    try:
                        actual_colors = [str(run.font.color.rgb) if run.font.color.type is not None else 'None' for run in body_runs]
                    except Exception:
                        actual_colors = ['Error']
                    print(f"FAIL: Component 4 — Body color values: {actual_colors}, expected all {TARGET_COLOR_STR}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in a given env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
