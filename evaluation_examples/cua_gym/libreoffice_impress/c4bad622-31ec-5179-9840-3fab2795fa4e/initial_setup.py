"""
Initial Setup: 5-slide financial report presentation - before underline/dark-red formatting task
Task ID: osworld_impress_underline_darkred_table_003
Domain: libreoffice_impress

Slide 4 has a title and 4 bullet points with NO underline and NO dark red (#8B0000) color.
The agent must apply underline and #8B0000 to all text on slide 4.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_underline_darkred_table_003'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Annual Financial Report 2024"
    slide1.placeholders[1].text = "Prepared by: Finance & Strategy Division"
    for para in slide1.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.bold = True
            run.font.size = Pt(36)
            run.font.color.rgb = RGBColor(0x1F, 0x39, 0x7C)

    # ---- Slide 2: Executive Summary ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Total Revenue: $4.82B (up 12% YoY)"
    tf2.add_paragraph().text = "Operating Income: $1.14B (margin 23.7%)"
    tf2.add_paragraph().text = "Net Profit: $874M, EPS: $3.42"
    tf2.add_paragraph().text = "Headcount grew from 18,200 to 19,650"
    for para in slide2.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1F, 0x39, 0x7C)

    # ---- Slide 3: Revenue Breakdown ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Revenue Breakdown by Segment"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "North America: $2.31B (+9%)"
    tf3.add_paragraph().text = "Europe: $1.28B (+15%)"
    tf3.add_paragraph().text = "Asia-Pacific: $0.97B (+22%)"
    tf3.add_paragraph().text = "Rest of World: $0.26B (+4%)"
    for para in slide3.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1F, 0x39, 0x7C)

    # ---- Slide 4: Cost Analysis ----
    # NOTE: NO underline, NO dark red (#8B0000) — agent must apply these
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Cost Analysis & Efficiency"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "COGS reduced by 3.2% through supplier renegotiation"
    p4b = tf4.add_paragraph()
    p4b.text = "R&D spend: $412M (8.5% of revenue)"
    p4c = tf4.add_paragraph()
    p4c.text = "SG&A: $654M — streamlined via automation initiatives"
    p4d = tf4.add_paragraph()
    p4d.text = "Capital expenditure: $298M focused on data centers"
    # Slide 4 title: plain black, not underlined
    for para in slide4.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1F, 0x39, 0x7C)
    # Slide 4 body: plain text, no underline, no dark red
    for para in slide4.placeholders[1].text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            run.font.underline = False

    # ---- Slide 5: Outlook ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "2025 Outlook & Strategic Priorities"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Revenue guidance: $5.3B–$5.5B (+10–14% YoY)"
    tf5.add_paragraph().text = "Target operating margin: 25%+"
    tf5.add_paragraph().text = "Planned M&A activity: 2 acquisitions in pipeline"
    tf5.add_paragraph().text = "Sustainability: net-zero carbon by 2030"
    for para in slide5.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1F, 0x39, 0x7C)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
