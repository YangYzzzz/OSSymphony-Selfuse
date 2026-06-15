"""
Initial Setup: Create Budget_Tracker presentation with a table on slide 3
Task ID: impress_tct_021
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_021'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_cell_text(cell, text, font_name="Liberation Sans", font_size=Pt(10),
                  bold=False, alignment=PP_ALIGN.LEFT, font_color=None):
    """Set cell text with explicit formatting."""
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    if font_color:
        run.font.color.rgb = font_color


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Budget Tracker 2025"
    slide1.placeholders[1].text = "Q1 Financial Overview - Greenfield Corp"

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Budget Overview"
    body2 = slide2.placeholders[1]
    tf2 = body2.text_frame
    tf2.text = "This presentation tracks departmental spending across Q1 2025."
    p2 = tf2.add_paragraph()
    p2.text = "Each department submitted monthly reports for January through March."
    p3 = tf2.add_paragraph()
    p3.text = "Total allocated budget: $1,250,000"

    # --- Slide 3: Table with budget data ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Add title text box
    from pptx.util import Emu
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Q1 Department Spending"
    run.font.name = "Liberation Sans"
    run.font.size = Pt(20)
    run.font.bold = True

    # Table: 8 rows x 3 columns
    rows, cols = 8, 3
    table_shape = slide3.shapes.add_table(
        rows, cols, Inches(1.5), Inches(1.2), Inches(6.5), Inches(4.0)
    )
    table = table_shape.table

    # Set column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(2.0)

    # Header row (row 0) - 10pt Liberation Sans, left-aligned, bold white on dark blue
    headers = ["Department", "Budget ($)", "Spent ($)"]
    header_bg = RGBColor(0x2E, 0x4A, 0x6E)
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        set_cell_text(cell, h, font_name="Liberation Sans", font_size=Pt(10),
                      bold=True, alignment=PP_ALIGN.LEFT,
                      font_color=RGBColor(0xFF, 0xFF, 0xFF))
        # Set header background
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = header_bg

    # Data rows (rows 1-7) - 10pt Liberation Sans, left-aligned
    data = [
        ["Engineering",    "320,000", "298,450"],
        ["Marketing",      "185,000", "172,300"],
        ["Sales",          "210,000", "195,680"],
        ["Human Resources", "95,000",  "88,120"],
        ["Operations",    "175,000", "163,900"],
        ["Finance",       "140,000", "131,250"],
        ["Research & Dev", "125,000", "119,800"],
    ]

    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            set_cell_text(cell, val, font_name="Liberation Sans", font_size=Pt(10),
                          bold=False, alignment=PP_ALIGN.LEFT)

    # --- Slide 4: Summary ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Key Takeaways"
    body4 = slide4.placeholders[1]
    tf4 = body4.text_frame
    tf4.text = "All departments stayed within budget allocation for Q1."
    p4a = tf4.add_paragraph()
    p4a.text = "Engineering had the highest absolute spending at $298,450."
    p4b = tf4.add_paragraph()
    p4b.text = "Total Q1 spending: $1,169,500 out of $1,250,000 allocated."

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
