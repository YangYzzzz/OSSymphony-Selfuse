"""
Reward Script: Set data cell text in table on slide 3 to right-aligned, 11pt Arial
Task ID: impress_tct_021
Domain: libreoffice_impress
Scoring:
  Component 1 (0.40): Font name changed to Arial in data rows (rows 1-7)
  Component 2 (0.35): Font size changed to 11pt (139700 EMU) in data rows
  Component 3 (0.25): Alignment changed to RIGHT in data rows AND header row unchanged
"""

import os
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_021'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    # Find the table on slide 3 (index 2)
    slide = prs.slides[2]
    table = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            break

    if table is None:
        print("FAIL: No table found on slide 3")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Table should be 8 rows x 3 cols
    num_rows = len(table.rows)
    num_cols = len(table.columns)
    if num_rows < 2 or num_cols < 1:
        print(f"FAIL: Table dimensions too small: {num_rows}x{num_cols}")
        print("REWARD: 0.0")
        return 0.0

    # Collect font/alignment info from data rows (rows 1 through num_rows-1)
    data_cells_total = 0
    arial_count = 0
    size_11_count = 0
    right_align_count = 0

    EXPECTED_SIZE = 139700  # 11pt in EMU

    for r in range(1, num_rows):
        for c in range(num_cols):
            cell = table.cell(r, c)
            for para in cell.text_frame.paragraphs:
                runs = [run for run in para.runs if (run.text or "").strip()]
                if not runs:
                    continue
                data_cells_total += 1

                # Check alignment (per paragraph)
                if para.alignment == PP_ALIGN.RIGHT:
                    right_align_count += 1

                # Check font properties (per run)
                all_arial = all(run.font.name == 'Arial' for run in runs)
                all_11pt = all(run.font.size == EXPECTED_SIZE for run in runs)

                if all_arial:
                    arial_count += 1
                if all_11pt:
                    size_11_count += 1

    if data_cells_total == 0:
        print("FAIL: No data cells with text found in rows 1+")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Found {data_cells_total} data cell paragraphs with text")

    # Component 1: Font name changed to Arial in data rows (0.40 points)
    try:
        arial_ratio = arial_count / data_cells_total
        if arial_ratio >= 0.95:
            print(f"PASS: Component 1 - Font Arial in {arial_count}/{data_cells_total} data cells (0.40 pts)")
            total_score += 0.40
        elif arial_ratio > 0:
            partial = round(0.40 * arial_ratio, 3)
            print(f"PARTIAL: Component 1 - Font Arial in {arial_count}/{data_cells_total} data cells ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 - No data cells have Arial font (0/{data_cells_total})")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Font size changed to 11pt in data rows (0.35 points)
    try:
        size_ratio = size_11_count / data_cells_total
        if size_ratio >= 0.95:
            print(f"PASS: Component 2 - Font size 11pt in {size_11_count}/{data_cells_total} data cells (0.35 pts)")
            total_score += 0.35
        elif size_ratio > 0:
            partial = round(0.35 * size_ratio, 3)
            print(f"PARTIAL: Component 2 - Font size 11pt in {size_11_count}/{data_cells_total} data cells ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 - No data cells have 11pt font size (0/{data_cells_total})")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Alignment changed to RIGHT in data rows AND header unchanged (0.25 points)
    # This is a compound check: data rows must be right-aligned, AND header must NOT be right-aligned
    try:
        # Check header is unchanged (not right-aligned, not Arial 11pt)
        header_unchanged = True
        header_issues = []
        for c in range(num_cols):
            cell = table.cell(0, c)
            for para in cell.text_frame.paragraphs:
                runs = [run for run in para.runs if (run.text or "").strip()]
                if not runs:
                    continue
                if para.alignment == PP_ALIGN.RIGHT:
                    header_unchanged = False
                    header_issues.append(f"Col {c}: alignment is RIGHT")
                for run in runs:
                    if run.font.name == 'Arial' and run.font.size == EXPECTED_SIZE:
                        header_unchanged = False
                        header_issues.append(f"Col {c}: font changed to Arial 11pt")

        align_ratio = right_align_count / data_cells_total

        if align_ratio >= 0.95 and header_unchanged:
            print(f"PASS: Component 3 - Right-aligned in {right_align_count}/{data_cells_total} data cells, header unchanged (0.25 pts)")
            total_score += 0.25
        elif align_ratio >= 0.95 and not header_unchanged:
            # Data rows correct but header was modified too - partial credit
            print(f"PARTIAL: Component 3 - Data right-aligned but header also changed: {'; '.join(header_issues)} (0.15 pts)")
            total_score += 0.15
        elif align_ratio > 0:
            partial = round(0.25 * align_ratio, 3)
            print(f"PARTIAL: Component 3 - Right-aligned in {right_align_count}/{data_cells_total} data cells ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 - No data cells are right-aligned (0/{data_cells_total})")
            if not header_unchanged:
                print(f"  Also: Header was modified: {'; '.join(header_issues)}")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
