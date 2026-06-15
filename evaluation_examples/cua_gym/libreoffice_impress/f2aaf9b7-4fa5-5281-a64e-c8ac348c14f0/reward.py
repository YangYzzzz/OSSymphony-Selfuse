"""
Reward Script: Standardize all text boxes to use 1.5 line spacing
Task ID: impstruct_024
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5 pts): High compliance rate (>=90% paragraphs at 1.5 spacing)
  Component 2 (0.3 pts): Perfect compliance (100% paragraphs at 1.5 spacing)
  Component 3 (0.2 pts): All 6 slides individually compliant (every slide's paragraphs at 1.5)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impstruct_024'
TARGET_SPACING = 1.5  # 150% proportional line spacing
TARGET_PCT_VAL = 150000  # XML spcPct val for 150%


def persist_app_state(domain):
    """Save any unsaved LibreOffice edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in ("libreoffice_calc", "libreoffice_writer", "libreoffice_impress"):
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_paragraph_spacing_info(prs):
    """
    Collect line spacing info for all non-empty paragraphs across all slides.
    Returns list of dicts with slide_idx, shape_name, para_idx, spacing, text.
    """
    from pptx.oxml.ns import qn
    results = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for pi, para in enumerate(shape.text_frame.paragraphs):
                text = para.text.strip()
                if not text:
                    continue
                # Get spacing from python-pptx API
                spacing = para.line_spacing
                # Also check XML for precise value
                pPr = para._p.find(qn('a:pPr'))
                xml_pct = None
                if pPr is not None:
                    lnSpc = pPr.find(qn('a:lnSpc'))
                    if lnSpc is not None:
                        spc_pct = lnSpc.find(qn('a:spcPct'))
                        if spc_pct is not None:
                            xml_pct = int(spc_pct.get('val', '0'))
                results.append({
                    'slide_idx': si,
                    'shape_name': shape.name,
                    'para_idx': pi,
                    'spacing': spacing,
                    'xml_pct': xml_pct,
                    'text': text[:50],
                })
    return results


def is_spacing_150(info):
    """Check if a paragraph has exactly 1.5 line spacing (150%)."""
    # Check XML percentage value first (most precise)
    if info['xml_pct'] is not None:
        return info['xml_pct'] == TARGET_PCT_VAL
    # Fallback to python-pptx API value
    if info['spacing'] is not None:
        return abs(info['spacing'] - TARGET_SPACING) < 0.01
    # None means default/inherited (not 1.5)
    return False


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file must have 6 slides
    num_slides = len(prs.slides)
    if num_slides != 6:
        print(f"PRECONDITION FAIL: Expected 6 slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    # Gather all paragraph spacing info
    all_paras = get_paragraph_spacing_info(prs)
    total_paras = len(all_paras)
    if total_paras == 0:
        print("PRECONDITION FAIL: No non-empty paragraphs found")
        print("REWARD: 0.0")
        return 0.0

    # Count paragraphs with correct 1.5 spacing
    correct_paras = [p for p in all_paras if is_spacing_150(p)]
    incorrect_paras = [p for p in all_paras if not is_spacing_150(p)]
    compliance_rate = len(correct_paras) / total_paras

    print(f"INFO: Total non-empty paragraphs: {total_paras}")
    print(f"INFO: Paragraphs with 1.5 spacing: {len(correct_paras)}")
    print(f"INFO: Compliance rate: {compliance_rate:.2%}")

    # Component 1: High compliance rate (0.5 points)
    # Requires >= 90% of paragraphs to have 1.5 spacing
    # Awards partial credit proportional to compliance above initial baseline
    try:
        if compliance_rate >= 0.9:
            print(f"PASS: Component 1 -- {compliance_rate:.0%} paragraphs at 1.5 spacing (>= 90%) (0.5 pts)")
            total_score += 0.5
        elif compliance_rate > 0.0:
            # Partial credit: scale 0 to 0.5 based on compliance rate
            # But only if compliance > what initial file would have (which is 0%)
            partial = round(compliance_rate * 0.5, 2)
            print(f"PARTIAL: Component 1 -- {compliance_rate:.0%} compliance, awarding {partial} pts")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- 0% compliance, no paragraphs at 1.5 spacing")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Perfect compliance (0.3 points)
    # ALL paragraphs must have exactly 1.5 spacing
    try:
        if compliance_rate == 1.0:
            print(f"PASS: Component 2 -- 100% compliance, all {total_paras} paragraphs at 1.5 (0.3 pts)")
            total_score += 0.3
        else:
            # List first few non-compliant paragraphs for debugging
            for p in incorrect_paras[:5]:
                print(f"FAIL: Component 2 -- S{p['slide_idx']+1} {p['shape_name']} P{p['para_idx']}: "
                      f"spacing={p['spacing']}, xml_pct={p['xml_pct']} | {p['text']}")
            if len(incorrect_paras) > 5:
                print(f"  ... and {len(incorrect_paras) - 5} more non-compliant paragraphs")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: All slides individually compliant (0.2 points)
    # Every slide must have all its paragraphs at 1.5 spacing
    try:
        slides_compliant = 0
        for si in range(num_slides):
            slide_paras = [p for p in all_paras if p['slide_idx'] == si]
            slide_correct = [p for p in slide_paras if is_spacing_150(p)]
            if len(slide_paras) > 0 and len(slide_correct) == len(slide_paras):
                slides_compliant += 1
            else:
                wrong_count = len(slide_paras) - len(slide_correct)
                print(f"DETAIL: Slide {si+1}: {wrong_count}/{len(slide_paras)} paragraphs NOT at 1.5")

        if slides_compliant == num_slides:
            print(f"PASS: Component 3 -- All {num_slides} slides individually compliant (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 3 -- Only {slides_compliant}/{num_slides} slides fully compliant")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
