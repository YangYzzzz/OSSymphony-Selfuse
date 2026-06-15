"""
Initial Setup: Standardize all text boxes to use 1.5 line spacing
Task ID: impstruct_024
Domain: libreoffice_impress

Creates a 6-slide strategy deck with inconsistent line spacing across text boxes.
Some paragraphs use single (1.0), some 1.15, some double (2.0), and some default (None).
None of them use 1.5 line spacing.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impstruct_024'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, paragraphs, line_spacings):
    """
    Add a textbox with multiple paragraphs, each with a specified line spacing.
    paragraphs: list of (text, font_size, bold, alignment, color) tuples
    line_spacings: list of float or None values for each paragraph
    """
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (para_info, spacing) in enumerate(zip(paragraphs, line_spacings)):
        text, font_size, bold, alignment, color = para_info
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = RGBColor(*color)
        if alignment:
            p.alignment = alignment

        if spacing is not None:
            p.line_spacing = spacing


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ========== SLIDE 1: Title Slide ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title
    add_textbox(slide1, 1.5, 0.8, 10, 1.5,
        [("Nextera Solutions — Q3 2025 Strategy Review", 36, True, PP_ALIGN.CENTER, (0x1B, 0x3A, 0x6B))],
        [1.0]  # single spacing
    )
    # Subtitle
    add_textbox(slide1, 2.5, 2.8, 8, 1.0,
        [("Prepared by the Office of Strategy & Growth", 18, False, PP_ALIGN.CENTER, (0x5A, 0x5A, 0x5A)),
         ("Confidential — Internal Use Only", 14, True, PP_ALIGN.CENTER, (0x99, 0x33, 0x33))],
        [2.0, 2.0]  # double spacing
    )
    # Date box
    add_textbox(slide1, 4.0, 4.5, 5, 0.6,
        [("August 15, 2025  |  Board Presentation", 14, False, PP_ALIGN.CENTER, (0x77, 0x77, 0x77))],
        [None]  # default/no explicit spacing
    )

    # ========== SLIDE 2: Market Analysis ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide2, 0.5, 0.3, 12, 1.0,
        [("Market Analysis & Competitive Landscape", 30, True, PP_ALIGN.LEFT, (0x1B, 0x3A, 0x6B))],
        [1.15]  # 1.15 spacing
    )
    add_textbox(slide2, 0.5, 1.5, 5.8, 5.0,
        [("Key Market Trends", 18, True, PP_ALIGN.LEFT, (0x2C, 0x2C, 0x2C)),
         ("The enterprise SaaS market grew 23% YoY, reaching $195B in total addressable market. Cloud-native architectures continue to dominate new deployments.", 13, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44)),
         ("Our primary segment (mid-market HR tech) expanded 31% as organizations accelerated digital transformation of workforce management.", 13, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44)),
         ("Regulatory changes in EU data sovereignty are creating both compliance costs and competitive moats for vendors with regional infrastructure.", 13, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44))],
        [1.15, 1.0, 1.0, 1.0]  # mixed: 1.15 header, single body
    )
    add_textbox(slide2, 6.8, 1.5, 5.8, 5.0,
        [("Competitor Movements", 18, True, PP_ALIGN.LEFT, (0x2C, 0x2C, 0x2C)),
         ("Workday acquired talent analytics startup Peakon for $700M, signaling aggressive expansion into our adjacent market.", 13, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44)),
         ("SAP SuccessFactors launched a free-tier offering targeting companies under 50 employees, pressuring our SMB pipeline.", 13, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44)),
         ("BambooHR raised $350M Series D at $3.2B valuation, suggesting accelerated R&D investment in automation features.", 13, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44))],
        [2.0, 2.0, 2.0, 2.0]  # all double spacing
    )

    # ========== SLIDE 3: Financial Performance ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide3, 0.5, 0.3, 12, 1.0,
        [("Q3 Financial Performance Summary", 30, True, PP_ALIGN.LEFT, (0x1B, 0x3A, 0x6B))],
        [1.0]  # single
    )
    add_textbox(slide3, 0.5, 1.5, 5.8, 2.5,
        [("Revenue Highlights", 18, True, PP_ALIGN.LEFT, (0x2C, 0x2C, 0x2C)),
         ("Total revenue: $42.8M (+18% YoY)", 14, False, PP_ALIGN.LEFT, (0x33, 0x33, 0x33)),
         ("Subscription revenue: $37.1M (+22% YoY)", 14, False, PP_ALIGN.LEFT, (0x33, 0x33, 0x33)),
         ("Professional services: $5.7M (-3% YoY)", 14, False, PP_ALIGN.LEFT, (0x33, 0x33, 0x33)),
         ("Net revenue retention: 118%", 14, False, PP_ALIGN.LEFT, (0x33, 0x33, 0x33))],
        [None, 1.15, 1.15, 1.15, 1.15]  # default header, 1.15 body
    )
    add_textbox(slide3, 0.5, 4.2, 5.8, 2.5,
        [("Cost Structure", 18, True, PP_ALIGN.LEFT, (0x2C, 0x2C, 0x2C)),
         ("COGS: $12.4M (29% of revenue, down from 32%)", 14, False, PP_ALIGN.LEFT, (0x33, 0x33, 0x33)),
         ("R&D spend: $9.8M (23% of revenue)", 14, False, PP_ALIGN.LEFT, (0x33, 0x33, 0x33)),
         ("S&M spend: $11.2M (26% of revenue)", 14, False, PP_ALIGN.LEFT, (0x33, 0x33, 0x33)),
         ("G&A: $4.1M (10% of revenue)", 14, False, PP_ALIGN.LEFT, (0x33, 0x33, 0x33))],
        [2.0, 1.0, 1.0, 1.0, 1.0]  # double header, single body
    )
    add_textbox(slide3, 6.8, 1.5, 5.8, 5.0,
        [("Profitability Metrics", 18, True, PP_ALIGN.LEFT, (0x2C, 0x2C, 0x2C)),
         ("Gross margin: 71% (up from 68% in Q2)", 14, False, PP_ALIGN.LEFT, (0x33, 0x33, 0x33)),
         ("Operating margin: 12.4% (first quarter above 12%)", 14, False, PP_ALIGN.LEFT, (0x33, 0x33, 0x33)),
         ("EBITDA: $7.3M (+35% YoY)", 14, False, PP_ALIGN.LEFT, (0x33, 0x33, 0x33)),
         ("Free cash flow: $5.1M (12% FCF margin)", 14, False, PP_ALIGN.LEFT, (0x33, 0x33, 0x33)),
         ("Cash & equivalents: $89.2M", 14, False, PP_ALIGN.LEFT, (0x33, 0x33, 0x33))],
        [1.15, 2.0, 2.0, 2.0, 2.0, 2.0]  # 1.15 header, double body
    )

    # ========== SLIDE 4: Product Roadmap ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide4, 0.5, 0.3, 12, 1.0,
        [("Product Roadmap — H2 2025 Priorities", 30, True, PP_ALIGN.LEFT, (0x1B, 0x3A, 0x6B))],
        [2.0]  # double
    )
    add_textbox(slide4, 0.5, 1.5, 3.8, 5.0,
        [("Platform Enhancements", 16, True, PP_ALIGN.LEFT, (0x0D, 0x6E, 0x6E)),
         ("Migrate remaining monolith services to microservices architecture by end of Q4.", 12, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44)),
         ("Implement real-time data streaming pipeline for workforce analytics dashboards.", 12, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44)),
         ("Achieve SOC 2 Type II certification for new cloud regions (Frankfurt, Sydney).", 12, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44))],
        [None, None, None, None]  # all default
    )
    add_textbox(slide4, 4.7, 1.5, 3.8, 5.0,
        [("AI & Automation", 16, True, PP_ALIGN.LEFT, (0x0D, 0x6E, 0x6E)),
         ("Launch predictive attrition model v2 with 85%+ accuracy target.", 12, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44)),
         ("Deploy conversational AI assistant for employee self-service HR queries.", 12, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44)),
         ("Integrate automated skills-gap analysis with learning management recommendations.", 12, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44))],
        [1.15, 1.15, 1.15, 1.15]  # all 1.15
    )
    add_textbox(slide4, 8.9, 1.5, 3.8, 5.0,
        [("User Experience", 16, True, PP_ALIGN.LEFT, (0x0D, 0x6E, 0x6E)),
         ("Redesign mobile app with native iOS/Android components for 40% faster load times.", 12, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44)),
         ("Implement progressive web app (PWA) for offline access to critical HR functions.", 12, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44)),
         ("Roll out accessibility compliance (WCAG 2.1 AA) across all customer-facing modules.", 12, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44))],
        [1.0, 2.0, 2.0, 2.0]  # single header, double body
    )

    # ========== SLIDE 5: Go-to-Market Strategy ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide5, 0.5, 0.3, 12, 1.0,
        [("Go-to-Market Strategy & Channel Expansion", 30, True, PP_ALIGN.LEFT, (0x1B, 0x3A, 0x6B))],
        [1.0]  # single
    )
    add_textbox(slide5, 0.5, 1.5, 5.8, 2.5,
        [("Direct Sales Initiatives", 18, True, PP_ALIGN.LEFT, (0x2C, 0x2C, 0x2C)),
         ("Expand enterprise sales team by 15 reps focused on Fortune 1000 accounts with 500+ employee count.", 13, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44)),
         ("Launch vertical-specific solution packages for healthcare, financial services, and manufacturing.", 13, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44)),
         ("Implement account-based marketing program targeting top 200 prospect accounts.", 13, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44))],
        [2.0, 1.15, 1.15, 1.15]  # double header, 1.15 body
    )
    add_textbox(slide5, 6.8, 1.5, 5.8, 2.5,
        [("Partner Channel", 18, True, PP_ALIGN.LEFT, (0x2C, 0x2C, 0x2C)),
         ("Onboard 25 new system integrator partners across EMEA and APAC regions.", 13, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44)),
         ("Launch certified implementation partner program with tiered incentive structure.", 13, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44)),
         ("Build co-selling motion with Salesforce and Microsoft Dynamics ecosystem partners.", 13, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44))],
        [None, 1.0, 1.0, 1.0]  # default header, single body
    )
    add_textbox(slide5, 0.5, 4.5, 12, 2.0,
        [("Target Metrics for H2 2025", 18, True, PP_ALIGN.LEFT, (0x2C, 0x2C, 0x2C)),
         ("New ARR: $18M  |  Logo Adds: 120  |  Partner-Sourced Revenue: 25%  |  Win Rate: 35%", 14, False, PP_ALIGN.CENTER, (0x1B, 0x3A, 0x6B))],
        [1.0, 2.0]  # mixed
    )

    # ========== SLIDE 6: Next Steps & Timeline ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide6, 0.5, 0.3, 12, 1.0,
        [("Next Steps & Key Milestones", 30, True, PP_ALIGN.LEFT, (0x1B, 0x3A, 0x6B))],
        [1.15]  # 1.15
    )
    add_textbox(slide6, 0.5, 1.5, 5.8, 4.5,
        [("Immediate Actions (30 Days)", 18, True, PP_ALIGN.LEFT, (0x2C, 0x2C, 0x2C)),
         ("Finalize Q4 hiring plan with VP Engineering and VP Sales — target 22 new hires.", 13, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44)),
         ("Complete SOC 2 readiness assessment with external auditor (Deloitte engagement).", 13, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44)),
         ("Submit Frankfurt data center provisioning request to AWS account team.", 13, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44)),
         ("Kick off account-based marketing pilot with first cohort of 50 target accounts.", 13, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44))],
        [1.0, 2.0, 2.0, 2.0, 2.0]  # single header, double body
    )
    add_textbox(slide6, 6.8, 1.5, 5.8, 4.5,
        [("Q4 Milestones", 18, True, PP_ALIGN.LEFT, (0x2C, 0x2C, 0x2C)),
         ("October: Beta launch of AI assistant to 10 pilot customers.", 13, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44)),
         ("November: Complete microservices migration for payroll and benefits modules.", 13, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44)),
         ("December: GA release of mobile app v3.0 with offline capabilities.", 13, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44)),
         ("December: Annual customer conference (NexConnect 2025) in Austin, TX.", 13, False, PP_ALIGN.LEFT, (0x44, 0x44, 0x44))],
        [1.15, None, None, None, None]  # 1.15 header, default body
    )
    add_textbox(slide6, 2.0, 6.2, 9, 0.8,
        [("Questions? Contact: strategy@nextera-solutions.com  |  Ext. 4200", 12, False, PP_ALIGN.CENTER, (0x88, 0x88, 0x88))],
        [1.0]  # single
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
