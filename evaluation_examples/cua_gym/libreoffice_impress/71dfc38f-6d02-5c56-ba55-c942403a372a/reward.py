"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 152 of my Impress file, Text Box 1 and Picture 1 are sitting too high. How do I snap both objects to the exact vertical center of the slide so they line up perfectly?
Generated: 2025-09-10 17:52:40
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def is_vertically_centered(shape, slide_height, tolerance_ratio=0.02):
    """Return True if the shape's vertical center is within the tolerance
    of the slide's vertical center."""
    shape_center_y = shape.top + shape.height / 2
    slide_center_y = slide_height / 2
    tolerance = slide_height * tolerance_ratio
    diff = abs(shape_center_y - slide_center_y)
    print(f"    Center diff for '{shape.name}': {diff} EMUs (tolerance {tolerance})")
    return diff <= tolerance


def verify_snap_to_center(file_path):
    """Verify that TextBox 1 and the first picture on slide 152 are both
    vertically centered. Progressive scoring is applied:
      • 0.3 points for finding TextBox 1
      • 0.3 points for finding a picture
      • 0.2 points if TextBox 1 is vertically centered
      • 0.2 points if the picture is vertically centered
    Must reach 1.0 for full credit.
    """
    max_score = 1.0
    score = 0.0

    # ------------------------------------------------------------------
    # 1. Load the presentation ------------------------------------------------
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Failed to load PPTX: {e}")
        print("REWARD: 0.0")
        return 0.0

    slide_index = 151  # Slide 152 (0-based index)
    if len(prs.slides) <= slide_index:
        print(f"✗ Presentation has only {len(prs.slides)} slides – need slide 152")
        print("REWARD: 0.0")
        return 0.0

    print(f"✓ Loaded presentation with {len(prs.slides)} slides")

    # ------------------------------------------------------------------
    # 2. Locate required shapes on slide 152 ----------------------------
    # ------------------------------------------------------------------
    slide = prs.slides[slide_index]

    textbox = None
    picture = None

    for shape in slide.shapes:
        if shape.name.strip() == "TextBox 1":
            textbox = shape
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and picture is None:
            picture = shape  # first picture on the slide

    # --- Scoring for locating shapes -----------------------------------
    if textbox is not None:
        print("✓ Found 'TextBox 1'")
        score += 0.3
    else:
        print("✗ 'TextBox 1' not found on slide 152")

    if picture is not None:
        print(f"✓ Found picture shape ('{picture.name}')")
        score += 0.3
    else:
        print("✗ No picture shape found on slide 152")

    # ------------------------------------------------------------------
    # 3. Verify vertical centering --------------------------------------
    # ------------------------------------------------------------------
    slide_height = prs.slide_height  # in EMUs

    if textbox is not None:
        if is_vertically_centered(textbox, slide_height):
            print("✓ 'TextBox 1' is vertically centered")
            score += 0.2
        else:
            print("✗ 'TextBox 1' is NOT vertically centered")

    if picture is not None:
        if is_vertically_centered(picture, slide_height):
            print("✓ Picture is vertically centered")
            score += 0.2
        else:
            print("✗ Picture is NOT vertically centered")

    # ------------------------------------------------------------------
    # 4. Final score -----------------------------------------------------
    # ------------------------------------------------------------------
    final_score = min(score, max_score)
    print(f"REWARD: {final_score}")
    return final_score


# ----------------------------------------------------------------------
# Execute verification when script is run directly ---------------------
# ----------------------------------------------------------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/on_slide_152_of_my_impress_file_text_box_1_and_picture_1_are_sitting_too_high_how_do_i_snap_both_obj_golden.pptx"
    verify_snap_to_center(FILE_PATH)
