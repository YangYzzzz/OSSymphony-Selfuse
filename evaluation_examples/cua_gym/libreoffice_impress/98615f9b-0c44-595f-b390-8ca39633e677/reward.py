"""
FINAL REWARD SCRIPT - SUCCESS
Task: I’m updating my Impress deck and want an agenda slide right after the cover (that’s slide 1). How do I insert a brand-new slide directly after it and set the layout to the exact "Title and Content" option?
Generated: 2025-09-10 16:08:36
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE

def verify_agenda_slide(file_path: str) -> float:
    """Verify that a presentation contains an agenda slide directly after the
    cover slide (i.e., as slide 2) using the exact "Title and Content" layout.

    Scoring (progressive):
      • 0.3 – Presentation has at least two slides (agenda slide position exists).
      • 0.5 – Slide 2’s layout name contains both "Title" and "Content"
               (indicating the required layout was chosen).
      • 0.2 – Slide 2 actually contains both a title placeholder and a body/content
               placeholder, proving the layout is functional.

    Returns a float in [0.0, 1.0]. Prints detailed diagnostics and the final
    reward as required ("REWARD: X.X").
    """
    max_score = 1.0
    score = 0.0

    print(f"Verifying presentation: {file_path}")

    # ---------- Prerequisite checks (NO POINTS AWARDED) ----------
    if not os.path.exists(file_path):
        print("✗ File not found")
        return 0.0
    if not file_path.lower().endswith('.pptx'):
        print("✗ File is not a .pptx file")
        return 0.0

    try:
        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        print(f"✓ Presentation loaded successfully with {slide_count} slides")
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        return 0.0

    # ---------- Requirement A: Deck has ≥ 2 slides ----------
    if slide_count >= 2:
        print("✓ Deck contains at least 2 slides")
        score += 0.3
    else:
        print("✗ Deck has fewer than 2 slides; agenda slide missing")
        return score  # Cannot meet other requirements

    # Work with slide 2 (index 1)
    slide2 = prs.slides[1]

    # ---------- Requirement B: Correct layout name ----------
    layout_name = slide2.slide_layout.name or ''
    print(f'Slide 2 layout name: "{layout_name}"')
    if 'title' in layout_name.lower() and 'content' in layout_name.lower():
        print('✓ Slide 2 uses "Title and Content" layout')
        score += 0.5
    else:
        print('✗ Slide 2 does not use the correct layout')

    # ---------- Requirement C: Required placeholders exist ----------
    title_found = False
    body_found = False

    for shape in slide2.shapes:
        if not shape.is_placeholder:
            continue
        ph_type = shape.placeholder_format.type
        if ph_type in (PP_PLACEHOLDER_TYPE.TITLE, PP_PLACEHOLDER_TYPE.CENTER_TITLE):
            title_found = True
        if ph_type in (PP_PLACEHOLDER_TYPE.BODY, PP_PLACEHOLDER_TYPE.OBJECT):
            body_found = True

    print(f'Title placeholder present: {title_found}; Content/Body placeholder present: {body_found}')
    if title_found and body_found:
        print('✓ Required placeholders are present on slide 2')
        score += 0.2
    else:
        print('✗ Required placeholders missing on slide 2')

    # ---------- Final score ----------
    final_score = min(score, max_score)
    print(f'Computed reward score: {final_score}')
    return final_score

# ---------------------- Script Entry Point ----------------------
if __name__ == '__main__':
    # Path provided by the task context
    pptx_path = (
        '/home/user/im_updating_my_impress_deck_and_want_an_agenda_slide_right_after_the_cover_thats_slide_1_how_do_i_in_golden.pptx'
    )
    reward = verify_agenda_slide(pptx_path)
    print(f'REWARD: {reward}')

