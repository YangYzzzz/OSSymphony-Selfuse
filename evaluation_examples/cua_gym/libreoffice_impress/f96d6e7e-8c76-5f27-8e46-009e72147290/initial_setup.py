"""
Initial Setup: Add borders to images in a portfolio presentation
Task ID: impress_fix_087
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_087'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
IMG_DIR = f'{WORKDIR}/portfolio_images'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_sample_images():
    """Create 10 distinct portfolio-style images with different colors and labels."""
    os.makedirs(IMG_DIR, exist_ok=True)

    image_specs = [
        ("Sunrise Over Mountains", (255, 140, 50), (40, 60, 100)),
        ("City Skyline at Dusk", (60, 80, 130), (200, 170, 100)),
        ("Coastal Lighthouse", (70, 150, 180), (220, 220, 200)),
        ("Autumn Forest Trail", (160, 100, 40), (80, 130, 50)),
        ("Desert Sand Dunes", (210, 180, 120), (100, 60, 30)),
        ("Northern Lights", (20, 40, 80), (50, 200, 120)),
        ("Lavender Fields", (120, 80, 160), (180, 140, 200)),
        ("Snow-Capped Peaks", (180, 200, 220), (100, 120, 150)),
        ("Tropical Waterfall", (30, 120, 60), (60, 180, 200)),
        ("Cherry Blossom Path", (220, 160, 180), (180, 80, 120)),
    ]

    paths = []
    for i, (label, bg_color, accent_color) in enumerate(image_specs):
        img = Image.new('RGB', (800, 600), bg_color)
        draw = ImageDraw.Draw(img)
        # Add some visual texture - rectangles and circles
        draw.rectangle([50, 50, 750, 550], outline=accent_color, width=3)
        draw.ellipse([300, 150, 500, 450], fill=accent_color)
        # Add label text
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 28)
        except (OSError, IOError):
            font = ImageFont.load_default()
        bbox = draw.textbbox((0, 0), label, font=font)
        tw = bbox[2] - bbox[0]
        draw.text(((800 - tw) // 2, 520), label, fill=(255, 255, 255), font=font)

        path = f'{IMG_DIR}/photo_{i+1:02d}.png'
        img.save(path)
        paths.append(path)

    return paths


def create_initial():
    image_paths = create_sample_images()

    prs = Presentation()
    # Standard 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_titles = [
        "Sunrise Over Mountains",
        "City Skyline at Dusk",
        "Coastal Lighthouse",
        "Autumn Forest Trail",
        "Desert Sand Dunes",
        "Northern Lights",
        "Lavender Fields",
        "Snow-Capped Peaks",
        "Tropical Waterfall",
        "Cherry Blossom Path",
    ]

    for i in range(10):
        # Use blank layout for clean portfolio look
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Add title text box at top
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = slide_titles[i]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        # Add image centered on slide, NO border
        img_left = Inches(2.5)
        img_top = Inches(1.5)
        img_width = Inches(8)
        img_height = Inches(5.5)
        pic = slide.shapes.add_picture(
            image_paths[i], img_left, img_top, img_width, img_height
        )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
