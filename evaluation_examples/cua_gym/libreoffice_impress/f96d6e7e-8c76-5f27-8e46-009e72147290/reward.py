"""
Reward Script: Add 2-point dark gray (#444444) border around all images (slides 1-10)
Task ID: impress_fix_087
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): Border exists on all 10 images (solidFill line element)
  Component 2 (0.35): Border width is 2pt (25400 EMU) on all 10 images
  Component 3 (0.35): Border color is #444444 on all 10 images
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_087'
EXPECTED_SLIDES = 10
EXPECTED_WIDTH_EMU = 25400  # 2 points = 25400 EMU (1 pt = 12700 EMU)
EXPECTED_COLOR = '444444'


def persist_app_state(domain: str):
    """Attempt to save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.oxml.ns import qn
    except ImportError as e:
        print(f"CRITICAL: Missing library: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have exactly 10 slides
    num_slides = len(prs.slides)
    if num_slides != EXPECTED_SLIDES:
        print(f"PRECONDITION FAIL: Expected {EXPECTED_SLIDES} slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    # Collect border info for each slide's image
    border_exists_count = 0
    width_correct_count = 0
    color_correct_count = 0

    for idx, slide in enumerate(prs.slides):
        slide_num = idx + 1
        # Find picture shape on this slide
        pic_shape = None
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pic_shape = shape
                break

        if pic_shape is None:
            print(f"  Slide {slide_num}: No picture found - skipping")
            continue

        # Check for line/border element
        sp = pic_shape._element
        ln = sp.find('.//' + qn('a:ln'))

        if ln is None:
            print(f"  Slide {slide_num}: No border (no <a:ln> element)")
            continue

        # Check solidFill exists
        sf = ln.find(qn('a:solidFill'))
        if sf is None:
            print(f"  Slide {slide_num}: Border element exists but no solidFill")
            continue

        # Border exists with solid fill
        border_exists_count += 1

        # Check width
        w_str = ln.get('w')
        if w_str is not None:
            try:
                w_val = int(w_str)
                # Allow small tolerance (within 5% of expected)
                if abs(w_val - EXPECTED_WIDTH_EMU) <= EXPECTED_WIDTH_EMU * 0.05:
                    width_correct_count += 1
                    print(f"  Slide {slide_num}: Width OK ({w_val} EMU)")
                else:
                    print(f"  Slide {slide_num}: Width WRONG - expected ~{EXPECTED_WIDTH_EMU}, got {w_val}")
            except ValueError:
                print(f"  Slide {slide_num}: Width not parseable: {w_str}")
        else:
            print(f"  Slide {slide_num}: No width attribute on <a:ln>")

        # Check color
        srgb = sf.find(qn('a:srgbClr'))
        if srgb is not None:
            clr = srgb.get('val', '').upper()
            if clr == EXPECTED_COLOR.upper():
                color_correct_count += 1
                print(f"  Slide {slide_num}: Color OK (#{clr})")
            else:
                print(f"  Slide {slide_num}: Color WRONG - expected #{EXPECTED_COLOR}, got #{clr}")
        else:
            print(f"  Slide {slide_num}: No srgbClr in solidFill")

    # Component 1: Border exists on all images (0.3 points, proportional)
    comp1_score = 0.0
    try:
        if border_exists_count > 0:
            comp1_score = 0.3 * (border_exists_count / EXPECTED_SLIDES)
            total_score += comp1_score
            print(f"\nPASS: Component 1 - Border exists on {border_exists_count}/{EXPECTED_SLIDES} images ({comp1_score:.3f} pts)")
        else:
            print(f"\nFAIL: Component 1 - No images have borders (0/{EXPECTED_SLIDES})")
    except Exception as e:
        print(f"\nERROR: Component 1 - {e}")

    # Component 2: Border width is 2pt on all images (0.35 points, proportional)
    comp2_score = 0.0
    try:
        if width_correct_count > 0:
            comp2_score = 0.35 * (width_correct_count / EXPECTED_SLIDES)
            total_score += comp2_score
            print(f"PASS: Component 2 - Correct width on {width_correct_count}/{EXPECTED_SLIDES} images ({comp2_score:.3f} pts)")
        else:
            print(f"FAIL: Component 2 - No images have correct 2pt width (0/{EXPECTED_SLIDES})")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Border color is #444444 on all images (0.35 points, proportional)
    comp3_score = 0.0
    try:
        if color_correct_count > 0:
            comp3_score = 0.35 * (color_correct_count / EXPECTED_SLIDES)
            total_score += comp3_score
            print(f"PASS: Component 3 - Correct color on {color_correct_count}/{EXPECTED_SLIDES} images ({comp3_score:.3f} pts)")
        else:
            print(f"FAIL: Component 3 - No images have correct #444444 color (0/{EXPECTED_SLIDES})")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
