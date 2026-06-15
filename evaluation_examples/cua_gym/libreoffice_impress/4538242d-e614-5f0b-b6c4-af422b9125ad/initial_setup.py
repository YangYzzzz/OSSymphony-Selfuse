"""
Initial Setup: Create a 15-slide confidential financial presentation
Task ID: impress_el_019
Domain: libreoffice_impress

Creates a presentation with sensitive financial data across 15 slides,
then opens it in LibreOffice Impress for the agent to work with.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_el_019'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a text box with formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tf


def add_table_slide(prs, title_text, headers, data, layout_idx=5):
    """Add a slide with a title and data table."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    # Title
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
                 title_text, font_size=24, bold=True,
                 color=RGBColor(0x1B, 0x3A, 0x5C))

    num_rows = len(data) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.2), Inches(9), Inches(0.4 * num_rows)
    )
    table = table_shape.table

    # Headers
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        from pptx.dml.color import RGBColor as RC
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    # Data rows
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(val)
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(10)

    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ===== Slide 1: Title Slide =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    add_text_box(slide1, Inches(1), Inches(1.5), Inches(8), Inches(1.5),
                 "Confidential Financial Report", font_size=36, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, Inches(1), Inches(3.2), Inches(8), Inches(0.8),
                 "Meridian Capital Holdings — FY2025 Annual Review",
                 font_size=20, color=RGBColor(0xCC, 0xCC, 0xCC),
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, Inches(1), Inches(4.5), Inches(8), Inches(0.6),
                 "STRICTLY CONFIDENTIAL — Board Members Only",
                 font_size=14, bold=True,
                 color=RGBColor(0xFF, 0x66, 0x66), alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, Inches(1), Inches(5.5), Inches(8), Inches(0.5),
                 "Prepared: March 28, 2025 | Version 3.1",
                 font_size=12, color=RGBColor(0x99, 0x99, 0x99),
                 alignment=PP_ALIGN.CENTER)

    # ===== Slide 2: Executive Summary =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide2, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
                 "Executive Summary", font_size=28, bold=True,
                 color=RGBColor(0x1B, 0x3A, 0x5C))
    summary_text = (
        "Total Revenue: $487.3M (+12.4% YoY)\n"
        "Net Income: $78.6M (+8.9% YoY)\n"
        "Operating Margin: 16.1% (vs 15.3% prior year)\n"
        "Cash & Equivalents: $142.8M\n"
        "Total Debt: $210.5M (Debt-to-Equity: 0.72)\n"
        "Earnings Per Share: $4.82 (+11.3% YoY)\n"
        "Dividend Yield: 2.8%"
    )
    add_text_box(slide2, Inches(0.5), Inches(1.3), Inches(9), Inches(5),
                 summary_text, font_size=16)

    # ===== Slide 3: Revenue Breakdown by Division =====
    add_table_slide(prs, "Revenue Breakdown by Division",
                    ["Division", "Q1 ($M)", "Q2 ($M)", "Q3 ($M)", "Q4 ($M)", "Total ($M)", "YoY %"],
                    [
                        ["Technology Solutions", "32.4", "35.1", "38.7", "41.2", "147.4", "+15.8%"],
                        ["Financial Services", "28.9", "30.2", "31.5", "33.8", "124.4", "+9.2%"],
                        ["Healthcare Analytics", "18.7", "20.1", "22.3", "24.6", "85.7", "+18.4%"],
                        ["Energy & Infrastructure", "15.3", "16.8", "17.2", "18.1", "67.4", "+7.1%"],
                        ["Consumer Products", "14.2", "15.0", "15.8", "17.4", "62.4", "+11.6%"],
                    ])

    # ===== Slide 4: Quarterly Revenue Trend =====
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide4, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
                 "Quarterly Revenue Trend (FY2023-FY2025)", font_size=24, bold=True,
                 color=RGBColor(0x1B, 0x3A, 0x5C))
    trend_data = (
        "FY2023 Q1: $92.1M  |  Q2: $98.4M  |  Q3: $101.7M  |  Q4: $108.2M  = $400.4M\n"
        "FY2024 Q1: $101.5M  |  Q2: $107.8M  |  Q3: $112.3M  |  Q4: $112.1M  = $433.7M\n"
        "FY2025 Q1: $109.5M  |  Q2: $117.2M  |  Q3: $125.5M  |  Q4: $135.1M  = $487.3M\n\n"
        "CAGR (3-Year): 10.3%"
    )
    add_text_box(slide4, Inches(0.5), Inches(1.3), Inches(9), Inches(5),
                 trend_data, font_size=14)

    # ===== Slide 5: Operating Expenses =====
    add_table_slide(prs, "Operating Expenses — FY2025",
                    ["Category", "Amount ($M)", "% of Revenue", "YoY Change"],
                    [
                        ["Cost of Goods Sold", "198.7", "40.8%", "+10.2%"],
                        ["Research & Development", "62.3", "12.8%", "+16.5%"],
                        ["Sales & Marketing", "48.9", "10.0%", "+8.7%"],
                        ["General & Administrative", "35.2", "7.2%", "+5.1%"],
                        ["Depreciation & Amortization", "22.8", "4.7%", "+3.9%"],
                        ["Restructuring Charges", "8.4", "1.7%", "-22.1%"],
                        ["Total Operating Expenses", "376.3", "77.2%", "+9.1%"],
                    ])

    # ===== Slide 6: Balance Sheet Highlights =====
    add_table_slide(prs, "Balance Sheet Highlights",
                    ["Item", "FY2025 ($M)", "FY2024 ($M)", "Change"],
                    [
                        ["Cash & Equivalents", "142.8", "118.4", "+20.6%"],
                        ["Accounts Receivable", "67.3", "58.9", "+14.3%"],
                        ["Total Current Assets", "248.6", "212.7", "+16.9%"],
                        ["Property & Equipment", "185.4", "172.1", "+7.7%"],
                        ["Goodwill & Intangibles", "312.8", "298.5", "+4.8%"],
                        ["Total Assets", "892.3", "821.6", "+8.6%"],
                        ["Total Liabilities", "502.1", "478.3", "+5.0%"],
                        ["Shareholders Equity", "390.2", "343.3", "+13.7%"],
                    ])

    # ===== Slide 7: Cash Flow Statement =====
    add_table_slide(prs, "Cash Flow Statement — FY2025",
                    ["Category", "Amount ($M)"],
                    [
                        ["Net Income", "78.6"],
                        ["Depreciation & Amortization", "22.8"],
                        ["Changes in Working Capital", "-12.4"],
                        ["Operating Cash Flow", "89.0"],
                        ["Capital Expenditures", "-34.2"],
                        ["Acquisitions", "-18.7"],
                        ["Investing Cash Flow", "-52.9"],
                        ["Debt Repayment", "-15.0"],
                        ["Dividends Paid", "-22.3"],
                        ["Financing Cash Flow", "-37.3"],
                        ["Net Change in Cash", "24.4"],
                    ])

    # ===== Slide 8: Debt Structure =====
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide8, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
                 "Debt Structure & Maturity Profile", font_size=24, bold=True,
                 color=RGBColor(0x1B, 0x3A, 0x5C))
    debt_text = (
        "Senior Secured Notes (5.25%, due 2028): $85.0M\n"
        "Revolving Credit Facility (SOFR+175bps, due 2027): $42.5M drawn of $100M\n"
        "Convertible Notes (3.75%, due 2029): $50.0M\n"
        "Term Loan B (SOFR+250bps, due 2030): $33.0M\n\n"
        "Total Debt: $210.5M\n"
        "Weighted Average Cost of Debt: 4.38%\n"
        "Interest Coverage Ratio: 5.2x\n"
        "Net Debt / EBITDA: 1.8x"
    )
    add_text_box(slide8, Inches(0.5), Inches(1.3), Inches(9), Inches(5),
                 debt_text, font_size=14)

    # ===== Slide 9: Segment Profitability =====
    add_table_slide(prs, "Segment Profitability Analysis",
                    ["Segment", "Revenue ($M)", "EBITDA ($M)", "Margin %", "ROIC %"],
                    [
                        ["Technology Solutions", "147.4", "32.4", "22.0%", "18.5%"],
                        ["Financial Services", "124.4", "24.9", "20.0%", "16.2%"],
                        ["Healthcare Analytics", "85.7", "18.9", "22.1%", "21.3%"],
                        ["Energy & Infrastructure", "67.4", "10.1", "15.0%", "11.8%"],
                        ["Consumer Products", "62.4", "8.1", "13.0%", "9.4%"],
                    ])

    # ===== Slide 10: Key Risks & Mitigation =====
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide10, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
                 "Key Risks & Mitigation Strategies", font_size=24, bold=True,
                 color=RGBColor(0x1B, 0x3A, 0x5C))
    risks_text = (
        "1. Market Concentration — Top 10 clients represent 34% of revenue\n"
        "   Mitigation: Diversification initiative targeting 25% by FY2027\n\n"
        "2. Regulatory Compliance — GDPR, SOX, HIPAA exposure across segments\n"
        "   Mitigation: $8.2M compliance technology investment in FY2026\n\n"
        "3. Talent Retention — 14.2% attrition rate in Technology Solutions\n"
        "   Mitigation: RSU refresh program and hybrid work policy expansion\n\n"
        "4. Currency Exposure — 22% of revenue denominated in EUR/GBP\n"
        "   Mitigation: Rolling 12-month hedging program (85% coverage)"
    )
    add_text_box(slide10, Inches(0.5), Inches(1.3), Inches(9), Inches(5.5),
                 risks_text, font_size=13)

    # ===== Slide 11: Capital Allocation Plan =====
    slide11 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide11, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
                 "Capital Allocation Plan — FY2026", font_size=24, bold=True,
                 color=RGBColor(0x1B, 0x3A, 0x5C))
    capalloc_text = (
        "Organic Growth (R&D + CapEx): $105.0M (48%)\n"
        "  - AI/ML Platform Enhancement: $28.0M\n"
        "  - Cloud Infrastructure Migration: $22.5M\n"
        "  - New Product Development: $18.0M\n"
        "  - Facility Modernization: $14.5M\n"
        "  - IT Systems Upgrade: $12.0M\n"
        "  - Other CapEx: $10.0M\n\n"
        "M&A / Strategic Investments: $55.0M (25%)\n"
        "Shareholder Returns (Dividends + Buybacks): $40.0M (18%)\n"
        "Debt Reduction: $20.0M (9%)\n\n"
        "Total Planned Deployment: $220.0M"
    )
    add_text_box(slide11, Inches(0.5), Inches(1.3), Inches(9), Inches(5.5),
                 capalloc_text, font_size=13)

    # ===== Slide 12: M&A Pipeline =====
    add_table_slide(prs, "M&A Pipeline — Active Targets",
                    ["Target", "Sector", "Est. Value ($M)", "Stage", "Synergies ($M)"],
                    [
                        ["Nexus Data Systems", "AI Analytics", "45-55", "Due Diligence", "8-12"],
                        ["CloudBridge Inc.", "Cloud Infra", "30-40", "LOI Signed", "6-9"],
                        ["MedTech Solutions", "Healthcare IT", "20-28", "Preliminary", "4-7"],
                        ["GreenGrid Energy", "Energy Tech", "15-22", "Screening", "3-5"],
                    ])

    # ===== Slide 13: Employee Compensation Data =====
    add_table_slide(prs, "Executive Compensation — FY2025 (CONFIDENTIAL)",
                    ["Name", "Title", "Base ($K)", "Bonus ($K)", "RSU ($K)", "Total ($K)"],
                    [
                        ["Victoria Harrington", "CEO", "850", "680", "2,400", "3,930"],
                        ["James Whitfield", "CFO", "625", "438", "1,600", "2,663"],
                        ["Priya Ramaswamy", "CTO", "600", "420", "1,800", "2,820"],
                        ["Robert Castellano", "COO", "575", "403", "1,200", "2,178"],
                        ["Elena Marchetti", "CLO", "525", "315", "900", "1,740"],
                        ["David Nakamura", "CHRO", "475", "285", "750", "1,510"],
                    ])

    # ===== Slide 14: FY2026 Financial Projections =====
    add_table_slide(prs, "FY2026 Financial Projections",
                    ["Metric", "Conservative", "Base Case", "Optimistic"],
                    [
                        ["Revenue ($M)", "512.0", "538.5", "572.0"],
                        ["Revenue Growth", "+5.1%", "+10.5%", "+17.4%"],
                        ["EBITDA ($M)", "88.6", "98.2", "112.5"],
                        ["EBITDA Margin", "17.3%", "18.2%", "19.7%"],
                        ["Net Income ($M)", "72.4", "84.8", "99.1"],
                        ["EPS", "$4.44", "$5.20", "$6.08"],
                        ["Free Cash Flow ($M)", "58.2", "72.5", "88.3"],
                    ])

    # ===== Slide 15: Closing / Disclaimer =====
    slide15 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = slide15.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    add_text_box(slide15, Inches(1), Inches(2), Inches(8), Inches(1),
                 "Thank You", font_size=36, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_text_box(slide15, Inches(1), Inches(3.5), Inches(8), Inches(2),
                 "This document contains material non-public information.\n"
                 "Distribution is restricted to authorized board members.\n"
                 "Any unauthorized disclosure may result in legal action.",
                 font_size=14, color=RGBColor(0xCC, 0xCC, 0xCC),
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide15, Inches(1), Inches(5.8), Inches(8), Inches(0.5),
                 "Meridian Capital Holdings | Investor Relations | ir@meridiancap.com",
                 font_size=11, color=RGBColor(0x88, 0x88, 0x88),
                 alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup — open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=3.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
