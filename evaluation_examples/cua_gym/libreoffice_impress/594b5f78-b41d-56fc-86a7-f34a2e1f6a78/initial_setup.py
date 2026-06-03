"""
Initial Setup: Configure notes master with logo and confidential footer
Task ID: impress_gf1_033
Domain: libreoffice_impress

Creates a 12-slide Internal Brief presentation with speaker notes,
plus a logo.png on the Desktop. Notes Master has no logo or footer.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont

WORKDIR = '/home/user'
TASK_ID = 'impress_gf1_033'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
LOGO_PATH = f'{WORKDIR}/Desktop/logo.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_logo():
    """Create a simple company logo PNG on the Desktop."""
    os.makedirs(f'{WORKDIR}/Desktop', exist_ok=True)
    img = Image.new('RGBA', (200, 80), (255, 255, 255, 0))
    draw = ImageDraw.Draw(img)
    # Draw a blue rectangle with rounded feel
    draw.rectangle([5, 5, 75, 75], fill=(0, 82, 165), outline=(0, 60, 130), width=2)
    # Draw "AC" text inside the rectangle
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 28)
    except IOError:
        font = ImageFont.load_default()
    draw.text((15, 20), "AC", fill=(255, 255, 255), font=font)
    # Draw company name next to it
    try:
        font_name = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 18)
    except IOError:
        font_name = ImageFont.load_default()
    draw.text((85, 25), "Apex\nCorp", fill=(0, 82, 165), font=font_name)
    img.save(LOGO_PATH)
    print(f'Logo created: {LOGO_PATH}')


def create_presentation():
    """Create a 12-slide Internal Brief presentation with speaker notes."""
    prs = Presentation()

    # Slide data: (layout_idx, title, content/subtitle, notes)
    slides_data = [
        (0, "Q2 2025 Internal Brief", "Apex Corporation — Strategy & Operations Update",
         "Welcome everyone. This quarterly brief covers our strategic priorities, financial highlights, and departmental updates for Q2 2025."),
        (1, "Agenda", "1. Financial Overview\n2. Product Roadmap\n3. Marketing Campaigns\n4. Engineering Updates\n5. HR & Talent\n6. Risk Assessment\n7. Q3 Outlook",
         "We'll cover seven key areas today. Please hold questions until the end of each section."),
        (1, "Financial Overview", "• Revenue: $12.4M (+18% YoY)\n• Gross Margin: 68.2%\n• Operating Expenses: $7.1M\n• Net Income: $1.82M\n• Cash Position: $24.6M",
         "Revenue growth was driven primarily by enterprise contracts signed in Q1. The margin improvement reflects our shift toward higher-value service tiers."),
        (1, "Product Roadmap — Apex Platform v3.2", "• Real-time analytics dashboard (shipped April 15)\n• Multi-tenant workspace support (Beta, June target)\n• API Gateway v2 with rate limiting (In Progress)\n• Mobile SDK for iOS and Android (Planning)\n• SSO integration with Okta/Azure AD (Complete)",
         "The analytics dashboard has already seen 340 daily active users since launch. Multi-tenant support is our highest priority for enterprise clients."),
        (1, "Marketing Campaigns", "• 'Accelerate 2025' webinar series — 2,400 registrations\n• LinkedIn thought leadership — 45K impressions/week\n• Partner co-marketing with TechForward Inc.\n• Industry conference sponsorships: SaaStr, Web Summit\n• Customer case study videos — 3 published, 5 in pipeline",
         "The webinar series exceeded our target by 60%. Sarah Chen's team has done exceptional work on the partner channel strategy."),
        (1, "Engineering Updates", "• Sprint velocity: 142 story points (up from 118)\n• Code coverage: 87.3% → 91.1%\n• P0 incidents: 0 in Q2\n• Infrastructure migration to Kubernetes: 78% complete\n• New hires onboarded: 4 senior engineers",
         "Zero P0 incidents is a significant achievement. The Kubernetes migration is on track for completion by end of Q3."),
        (1, "HR & Talent Pipeline", "• Current headcount: 156 (target: 170 by EOY)\n• Open positions: 14 across Engineering, Sales, CS\n• Employee satisfaction score: 4.3/5.0\n• Voluntary turnover rate: 8.2% (industry avg: 13.5%)\n• DEI initiative: Mentorship program launched",
         "We're ahead of industry benchmarks on retention. The mentorship program pairs junior team members with senior leaders across departments."),
        (1, "Risk Assessment", "• Supply chain delays affecting hardware partnerships\n• Regulatory changes in EU data residency (GDPR updates)\n• Competitive pressure from NovaTech's new pricing model\n• Cybersecurity: Completed SOC 2 Type II audit\n• Key person dependency in Platform Architecture team",
         "The SOC 2 audit completion strengthens our enterprise sales position. We need to address the key person risk in Platform Architecture."),
        (1, "Customer Success Metrics", "• NPS Score: 72 (up from 64)\n• Customer retention rate: 94.8%\n• Average response time: 2.1 hours\n• Support tickets resolved same-day: 88%\n• Upsell revenue: $1.2M",
         "NPS improvement reflects the investment in our support infrastructure. Marcus Johnson's team has been instrumental in driving upsell conversations."),
        (1, "Regional Performance", "• North America: $7.8M (63% of revenue)\n• Europe: $3.1M (25%, +22% YoY)\n• Asia-Pacific: $1.5M (12%, +35% YoY)\n• New markets: Middle East pilot with 3 accounts\n• Partner channel contribution: 18% of new ARR",
         "Asia-Pacific growth is accelerating faster than projected. The Middle East pilot could become a significant revenue stream by 2026."),
        (1, "Q3 Strategic Priorities", "1. Launch multi-tenant workspace support\n2. Complete Kubernetes migration\n3. Expand APAC sales team (+3 headcount)\n4. Finalize SOC 2 Type II certification marketing\n5. Begin Series C fundraising preparation",
         "These five priorities were aligned with the executive team last week. Series C preparation needs to begin now to target a Q1 2026 raise."),
        (0, "Thank You", "Questions & Discussion\nContact: strategy@apexcorp.com",
         "Thank you all for attending. Please reach out to your department heads for detailed breakdowns. Next quarterly brief is scheduled for October 8."),
    ]

    for layout_idx, title, content, notes in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        if slide.shapes.title:
            slide.shapes.title.text = title
        # Set content in the appropriate placeholder
        if layout_idx == 0:
            # Title slide: subtitle is placeholder[1]
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = content
        else:
            # Title+Content: content body is placeholder[1]
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = content
        # Add speaker notes
        slide.notes_slide.notes_text_frame.text = notes

    prs.save(OUTPUT)
    print(f'Presentation created: {OUTPUT} ({len(prs.slides)} slides)')


def create_initial():
    create_logo()
    create_presentation()

    # GUI-ready: open presentation in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
