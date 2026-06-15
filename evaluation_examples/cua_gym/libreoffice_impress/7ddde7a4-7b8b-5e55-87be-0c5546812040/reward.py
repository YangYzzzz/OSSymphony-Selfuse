"""
Reward Script: Set presentation to read-only mode and save to Desktop
Task ID: impress_fix_091
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): File exists at ~/Desktop/Final_Protected.pptx
  Component 2 (0.4): readOnlyRecommended="1" attribute set in presentation.xml
  Component 3 (0.3): Content integrity - 15 slides preserved with original content
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_091'
TARGET_PATH = os.path.join(WORKDIR, 'Desktop', 'Final_Protected.pptx')


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Component 1: File exists at ~/Desktop/Final_Protected.pptx (0.3 points)
    # This is a task-introduced change: the file must be saved to a NEW location
    try:
        if os.path.isfile(file_path):
            # Verify it's a valid ZIP (pptx format)
            if zipfile.is_zipfile(file_path):
                print(f"PASS: Component 1 - File exists at {file_path} and is valid pptx (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 1 - File exists but is not a valid pptx")
        else:
            print(f"FAIL: Component 1 - File not found at {file_path}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: readOnlyRecommended="1" in presentation.xml (0.4 points)
    # This is the core task change - setting read-only recommendation
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            if 'ppt/presentation.xml' in zf.namelist():
                with zf.open('ppt/presentation.xml') as f:
                    content = f.read().decode('utf-8')
                    root = ET.fromstring(content)
                    # Check for readOnlyRecommended attribute on the root p:presentation element
                    read_only_val = root.attrib.get('readOnlyRecommended', None)
                    if read_only_val == '1':
                        print(f"PASS: Component 2 - readOnlyRecommended='1' is set (0.4 pts)")
                        total_score += 0.4
                    elif read_only_val is not None:
                        print(f"FAIL: Component 2 - readOnlyRecommended found but value is '{read_only_val}', expected '1'")
                    else:
                        # Also check raw string in case namespace issues
                        if 'readOnlyRecommended="1"' in content or "readOnlyRecommended='1'" in content:
                            print(f"PASS: Component 2 - readOnlyRecommended='1' found via string search (0.4 pts)")
                            total_score += 0.4
                        else:
                            print(f"FAIL: Component 2 - readOnlyRecommended attribute not found in presentation.xml")
            else:
                print(f"FAIL: Component 2 - ppt/presentation.xml not found in archive")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Content integrity - 15 slides preserved (0.3 points)
    # We verify the file has 15 slides and key content from original presentation
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slide_count = len(prs.slides)

        if slide_count == 15:
            # Check that key slides have expected content
            slide1_text = ""
            for shape in prs.slides[0].shapes:
                if shape.has_text_frame:
                    slide1_text += shape.text_frame.text

            if "Q4 2025 Strategic Review" in slide1_text and "Meridian" in slide1_text:
                print(f"PASS: Component 3 - 15 slides present with correct content (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 3 - 15 slides present but title content mismatch: '{slide1_text[:80]}'")
        else:
            print(f"FAIL: Component 3 - Expected 15 slides, found {slide_count}")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
if not os.path.exists(TARGET_PATH):
    print(f"File not found: {TARGET_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(TARGET_PATH)
