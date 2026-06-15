"""
Initial Setup: Create an editable 15-slide business presentation Final_Draft.pptx
Task ID: impress_fix_091
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_091'
OUTPUT = f'{WORKDIR}/Final_Draft.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_shape(shape, text, font_size=18, bold=False, color=None, alignment=None):
    """Helper to set text properties on a shape."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Q4 2025 Strategic Review"
    slide.placeholders[1].text = "Meridian Consulting Group\nPrepared by Sarah Chen, VP Strategy"

    # --- Slide 2: Agenda ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Agenda"
    body = slide.placeholders[1].text_frame
    body.text = "1. Executive Summary"
    items = [
        "2. Financial Performance Overview",
        "3. Market Analysis & Competitive Landscape",
        "4. Client Portfolio Review",
        "5. Operational Efficiency Metrics",
        "6. Technology Infrastructure Update",
        "7. Talent & Workforce Planning",
        "8. Risk Assessment",
        "9. Strategic Initiatives for Q1 2026",
        "10. Q&A and Next Steps",
    ]
    for item in items:
        p = body.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 3: Executive Summary ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    body = slide.placeholders[1].text_frame
    body.text = "Revenue grew 18% YoY to $47.3M, exceeding target by $2.1M"
    for txt in [
        "Client retention rate improved to 94.2%, up from 91.8%",
        "Successfully launched 3 new service verticals",
        "Expanded into APAC region with Singapore office",
        "Net promoter score increased to 72 from 65",
    ]:
        p = body.add_paragraph()
        p.text = txt

    # --- Slide 4: Revenue Breakdown ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Revenue Breakdown by Segment"
    body = slide.placeholders[1].text_frame
    body.text = "Management Consulting: $18.9M (40%)"
    for txt in [
        "Digital Transformation: $12.3M (26%)",
        "Risk & Compliance: $9.5M (20%)",
        "Human Capital: $4.7M (10%)",
        "Other Services: $1.9M (4%)",
    ]:
        p = body.add_paragraph()
        p.text = txt

    # --- Slide 5: Client Portfolio ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Top Client Accounts"
    # Add a table
    table_shape = slide.shapes.add_table(6, 4, Inches(0.5), Inches(1.8), Inches(9), Inches(3.5))
    table = table_shape.table
    headers = ["Client", "Industry", "Annual Revenue", "Status"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
    clients = [
        ["Apex Financial Corp", "Banking", "$4.2M", "Active"],
        ["Northwind Pharma", "Healthcare", "$3.8M", "Active"],
        ["Cobalt Energy Ltd", "Energy", "$2.9M", "Renewal Due"],
        ["Sterling Retail Group", "Retail", "$2.4M", "Active"],
        ["Pinnacle Tech Solutions", "Technology", "$2.1M", "Expanding"],
    ]
    for r, row_data in enumerate(clients, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 6: Market Analysis ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Market Analysis"
    body = slide.placeholders[1].text_frame
    body.text = "Total addressable market grew to $285B globally"
    for txt in [
        "AI/ML consulting demand increased 42% year-over-year",
        "ESG advisory services emerging as fastest-growing segment",
        "Competitive pressure from Big 4 remains steady",
        "Mid-market firms consolidating at accelerated pace",
    ]:
        p = body.add_paragraph()
        p.text = txt

    # --- Slide 7: Competitive Landscape ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Competitive Positioning"
    body = slide.placeholders[1].text_frame
    body.text = "Meridian ranked #3 in mid-market consulting (Forrester)"
    for txt in [
        "Differentiated through proprietary analytics platform",
        "Higher client satisfaction scores than top 2 competitors",
        "Need to strengthen brand presence in EMEA markets",
    ]:
        p = body.add_paragraph()
        p.text = txt

    # --- Slide 8: Operational Efficiency ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Operational Efficiency Metrics"
    body = slide.placeholders[1].text_frame
    body.text = "Utilization rate: 78.4% (target 80%)"
    for txt in [
        "Average project margin: 34.2% (up from 31.8%)",
        "Project delivery on-time rate: 91.7%",
        "Proposal win rate: 38.5% (industry avg: 32%)",
        "Average engagement duration: 8.3 months",
    ]:
        p = body.add_paragraph()
        p.text = txt

    # --- Slide 9: Technology Update ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Technology Infrastructure"
    body = slide.placeholders[1].text_frame
    body.text = "Completed migration to Azure cloud platform"
    for txt in [
        "Deployed Meridian Analytics Suite v3.2 to all consultants",
        "CRM upgrade to Salesforce Lightning completed",
        "Cybersecurity audit: zero critical findings",
        "IT spend reduced 12% through vendor consolidation",
    ]:
        p = body.add_paragraph()
        p.text = txt

    # --- Slide 10: Talent Overview ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Talent & Workforce"
    body = slide.placeholders[1].text_frame
    body.text = "Total headcount: 312 (up from 278)"
    for txt in [
        "Voluntary turnover: 14.2% (industry avg 18.6%)",
        "Senior hires: 8 Partners, 15 Senior Managers",
        "Diversity index improved to 0.68 from 0.61",
        "Employee engagement score: 4.2/5.0",
    ]:
        p = body.add_paragraph()
        p.text = txt

    # --- Slide 11: Risk Assessment ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Risk Assessment"
    body = slide.placeholders[1].text_frame
    body.text = "Economic downturn risk: MEDIUM - diversified client base mitigates"
    for txt in [
        "Talent retention risk: LOW - competitive compensation packages",
        "Technology disruption risk: MEDIUM - AI investment ongoing",
        "Regulatory risk: LOW - compliance framework robust",
        "Concentration risk: MEDIUM - top 5 clients = 33% revenue",
    ]:
        p = body.add_paragraph()
        p.text = txt

    # --- Slide 12: Strategic Initiatives ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Strategic Initiatives for Q1 2026"
    body = slide.placeholders[1].text_frame
    body.text = "Launch AI Advisory Practice (Budget: $2.5M)"
    for txt in [
        "Expand Singapore office to 25 consultants",
        "Acquire boutique ESG firm (target: GreenPath Advisory)",
        "Roll out Meridian Academy internal training platform",
        "Establish strategic partnership with TechVault Inc.",
    ]:
        p = body.add_paragraph()
        p.text = txt

    # --- Slide 13: Financial Projections ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Q1 2026 Financial Projections"
    body = slide.placeholders[1].text_frame
    body.text = "Projected Revenue: $13.2M (+15% QoQ)"
    for txt in [
        "Target Operating Margin: 22.5%",
        "Capital Expenditure: $1.8M (technology & facilities)",
        "Expected New Clients: 12-15 engagements",
        "Pipeline Value: $42.7M (weighted)",
    ]:
        p = body.add_paragraph()
        p.text = txt

    # --- Slide 14: Timeline ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Implementation Timeline"
    body = slide.placeholders[1].text_frame
    body.text = "January 2026: AI Practice launch & team onboarding"
    for txt in [
        "February 2026: ESG acquisition due diligence complete",
        "March 2026: Meridian Academy beta launch",
        "April 2026: Singapore expansion phase 2",
        "May 2026: Mid-year strategic review checkpoint",
    ]:
        p = body.add_paragraph()
        p.text = txt

    # --- Slide 15: Q&A ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Questions & Discussion"
    slide.placeholders[1].text = "Thank you for your attention\nContact: s.chen@meridianconsulting.com"

    # Save the presentation
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
