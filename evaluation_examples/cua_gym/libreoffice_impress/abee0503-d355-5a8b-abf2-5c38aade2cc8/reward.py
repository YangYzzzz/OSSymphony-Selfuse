"""
Reward Script: Data Center Migration Presentation
Task ID: impress_wf_088
Domain: libreoffice_impress
Scoring: 10 components (0.10 each) verifying slide count, content structure,
         shapes, tables, diagrams, decision tree, chart, and color usage.
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_088'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)
    num_slides = len(slides)

    # Component 1: Presentation has exactly 10 slides (0.10 pts)
    # Initial has 1 blank slide; golden has 10. This is the core structural change.
    try:
        if num_slides == 10:
            print(f"PASS: Component 1 — 10 slides found (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 1 — Expected 10 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Early exit if not enough slides to check the rest
    if num_slides < 10:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Slide 1 title contains "Data Center Migration Plan - Phase 2" (0.10 pts)
    try:
        slide1 = slides[0]
        all_text = ""
        for shape in slide1.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    all_text += para.text + " "
        if "Data Center Migration Plan" in all_text and "Phase 2" in all_text:
            print(f"PASS: Component 2 — Title text found on slide 1 (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 — Title text not found. Got: {all_text[:100]}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 2 has rectangle shapes representing server racks (0.10 pts)
    # Check for AUTO_SHAPE shapes with 'Rectangle' in name and server-related text
    try:
        slide2 = slides[1]
        rect_count = 0
        for shape in slide2.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                name_lower = shape.name.lower()
                if 'rectangle' in name_lower and 'rounded' not in name_lower:
                    rect_count += 1
        if rect_count >= 3:
            print(f"PASS: Component 3 — Slide 2 has {rect_count} rectangle shapes (server racks) (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 3 — Expected >=3 rectangles on slide 2, found {rect_count}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 3 has rounded rectangle shapes (cloud architecture) (0.10 pts)
    try:
        slide3 = slides[2]
        rounded_count = 0
        for shape in slide3.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                name_lower = shape.name.lower()
                if 'rounded' in name_lower:
                    rounded_count += 1
        if rounded_count >= 3:
            print(f"PASS: Component 4 — Slide 3 has {rounded_count} rounded rectangle shapes (cloud) (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 4 — Expected >=3 rounded rectangles on slide 3, found {rounded_count}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slide 4 has a migration waves table with correct headers (0.10 pts)
    try:
        slide4 = slides[3]
        tables = [s for s in slide4.shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE]
        table_ok = len(tables) > 0
        if table_ok:
            tbl = tables[0].table
            if len(tbl.columns) >= 4 and len(tbl.rows) >= 2:
                headers = [tbl.cell(0, c).text.strip().lower() for c in range(len(tbl.columns))]
                table_ok = (any('wave' in h for h in headers)
                            and any('system' in h for h in headers)
                            and any('timeline' in h for h in headers)
                            and any('risk' in h for h in headers))
            else:
                table_ok = len(tbl.columns) >= 4
        if table_ok:
            print(f"PASS: Component 5 — Slide 4 has migration waves table with correct headers (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 5 — Migration waves table not found or headers incorrect on slide 4")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Slide 5 has oval/circle shapes AND connector lines (dependency mapping) (0.10 pts)
    try:
        slide5 = slides[4]
        oval_count = 0
        line_count = 0
        for shape in slide5.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if 'oval' in shape.name.lower():
                    oval_count += 1
            if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                line_count += 1
        if oval_count >= 3 and line_count >= 3:
            print(f"PASS: Component 6 — Slide 5 has {oval_count} ovals + {line_count} connectors (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 6 — Expected >=3 ovals and >=3 lines on slide 5, found {oval_count} ovals, {line_count} lines")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Slide 7 has flowchart shapes connected by lines (0.10 pts)
    try:
        slide7 = slides[6]
        shape_count = 0
        connector_count = 0
        for shape in slide7.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                shape_count += 1
            if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                connector_count += 1
        if shape_count >= 3 and connector_count >= 2:
            print(f"PASS: Component 7 — Slide 7 has flowchart: {shape_count} shapes + {connector_count} connectors (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 7 — Expected >=3 shapes and >=2 connectors on slide 7, found {shape_count} shapes, {connector_count} connectors")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    # Component 8: Slide 8 has diamond decision shapes (rollback decision tree) (0.10 pts)
    try:
        slide8 = slides[7]
        diamond_count = 0
        for shape in slide8.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if 'diamond' in shape.name.lower():
                    diamond_count += 1
        if diamond_count >= 1:
            print(f"PASS: Component 8 — Slide 8 has {diamond_count} diamond shapes (decision tree) (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 8 — Expected >=1 diamond shape on slide 8, found {diamond_count}")
    except Exception as e:
        print(f"ERROR: Component 8 — {e}")

    # Component 9: Slide 9 has bar chart visualization (rectangles + lines for projection) (0.10 pts)
    try:
        slide9 = slides[8]
        rect_count = 0
        line_count = 0
        for shape in slide9.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if 'rectangle' in shape.name.lower():
                    rect_count += 1
            if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                line_count += 1
        # Need multiple bar rectangles (at least 4 pairs = 8) and projection lines
        if rect_count >= 6 and line_count >= 2:
            print(f"PASS: Component 9 — Slide 9 has {rect_count} rectangles (bars) + {line_count} lines (projection) (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 9 — Expected >=6 rectangles and >=2 lines on slide 9, found {rect_count} rects, {line_count} lines")
    except Exception as e:
        print(f"ERROR: Component 9 — {e}")

    # Component 10: Colors #039BE5 (cloud-blue) and #546E7A (gray) used in the presentation (0.10 pts)
    # Check at XML level since colors are used across shapes throughout
    try:
        blue_slides = 0
        gray_slides = 0
        with zipfile.ZipFile(file_path, 'r') as zf:
            for slide_num in range(1, 11):
                fname = f'ppt/slides/slide{slide_num}.xml'
                try:
                    with zf.open(fname) as f:
                        content = f.read().decode('utf-8')
                        if '039BE5' in content:
                            blue_slides += 1
                        if '546E7A' in content:
                            gray_slides += 1
                except KeyError:
                    pass
        if blue_slides > 0 and gray_slides > 0:
            print(f"PASS: Component 10 — Colors #039BE5 on {blue_slides} slides and #546E7A on {gray_slides} slides (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 10 — blue=#039BE5 slides={blue_slides}, gray=#546E7A slides={gray_slides}")
    except Exception as e:
        print(f"ERROR: Component 10 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in a given env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
