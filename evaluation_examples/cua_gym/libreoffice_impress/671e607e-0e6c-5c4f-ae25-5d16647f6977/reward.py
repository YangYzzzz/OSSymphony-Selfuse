"""
FINAL REWARD SCRIPT - SUCCESS
Task: I just finished formatting Slide 4, and its heading is in the hex color #1E90FF. Could you drop in a fresh slide right after that and give the new title the exact text "Online Shopping", styled with the same #1E90FF font color?
Generated: 2025-09-10 13:37:56
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
from pptx import Presentation


def verify_online_shopping_slide(file_path: str) -> float:
    """Verify that a new slide (immediately after slide 4) contains the title
    'Online Shopping' in the exact font colour #1E90FF (RGB 30,144,255).

    Returns a progressive score between 0.0 and 1.0.
    """

    print(f"Verifying presentation: {file_path}\n")
    total_score = 0.0                # progressive score
    MAX_SCORE = 1.0                  # cap
    expected_hex = "1E90FF"          # expected RGB in hex (uppercase, no '#')

    # ------------------------------------------------------------------
    # 1. Load the presentation – prerequisite, no points just for loading
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File does not exist – task failed")
        return 0.0

    try:
        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        print(f"✓ Presentation loaded successfully with {slide_count} slides\n")
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        return 0.0

    # ---------------------------------------------------------------
    # 2. Basic structural check: at least 5 slides (original 4 + new)
    # ---------------------------------------------------------------
    if slide_count >= 5:
        print("✓ Presentation contains at least 5 slides – new slide appears to be present")
        total_score += 0.2  # 20 % for correct slide addition
    else:
        print("✗ Fewer than 5 slides – new slide missing")
        # Cannot continue reliably, return whatever score accrued (0.0 here)
        return total_score

    # ---------------------------------------------------------------
    # 3. Extract the relevant slides
    #    • slide 4 (index 3) – previous slide (for context only)
    #    • slide 5 (index 4) – expected new slide
    # ---------------------------------------------------------------
    try:
        slide5 = prs.slides[4]  # zero-based index
    except IndexError:
        print("✗ Slide 5 not found – cannot verify new slide")
        return total_score

    # Helper generator: yield (text, run) for every text run on a slide
    def iter_text_runs(slide):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    yield run.text, run

    # ---------------------------------------------------------------
    # 4. Requirement A – Title text must be exactly 'Online Shopping'
    # ---------------------------------------------------------------
    online_text_found = False
    colour_correct        = False

    for text, run in iter_text_runs(slide5):
        if text.strip().lower() == "online shopping":
            online_text_found = True
            # Check run colour (may raise if colour not set as RGB)
            try:
                rgb = run.font.color.rgb  # returns RGBColor or raises
            except Exception:
                rgb = None
            if rgb is not None and str(rgb).upper() == expected_hex:
                colour_correct = True
            # No break – continue scanning to catch multiple occurrences

    if online_text_found:
        print("✓ Found title text 'Online Shopping' on slide 5")
        total_score += 0.4  # 40 % for correct text
    else:
        print("✗ Title text 'Online Shopping' NOT found on slide 5")

    # ---------------------------------------------------------------
    # 5. Requirement B – Font colour must be exactly #1E90FF
    # ---------------------------------------------------------------
    if colour_correct:
        print("✓ Font colour for 'Online Shopping' is exactly #1E90FF")
        total_score += 0.4  # 40 % for correct colour
    else:
        if online_text_found:
            print("✗ 'Online Shopping' found but font colour is incorrect or not set to #1E90FF")
        else:
            print("(Skipping colour check – title text was missing)")

    # ---------------------------------------------------------------
    # 6. Final score calculation
    # ---------------------------------------------------------------
    final_score = min(total_score, MAX_SCORE)
    print(f"\nTotal score: {final_score}/{MAX_SCORE}")
    return final_score


# ----------------------------
# Driver code (invoked by VM)
# ----------------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/i_just_finished_formatting_slide_4_and_its_heading_is_in_the_hex_color_1e90ff_could_you_drop_in_a_fr_golden.pptx"
    reward = verify_online_shopping_slide(FILE_PATH)
    print(f"REWARD: {reward}")

