"""
Initial Setup: Add a 2-column by 4-row table on slide 6
Task ID: impress_tct_031
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_031'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_body(slide, title_text, body_lines):
    """Helper to populate title and content slide."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    # Find the body placeholder (index 1)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.text = body_lines[0]
            for line in body_lines[1:]:
                p = tf.add_paragraph()
                p.text = line
                p.level = 0
            break


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Decision Matrix"
    slide1.placeholders[1].text = "Strategic Evaluation Framework\nQ2 2026 Planning"

    # --- Slide 2: Background & Context ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide2, "Background & Context", [
        "Our team has been evaluating three vendor proposals for the new CRM platform",
        "Current system contract expires in September 2026",
        "Budget allocation approved at $450K for implementation",
        "Migration timeline target: 12 weeks post-selection",
    ])

    # --- Slide 3: Key Criteria ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide3, "Key Criteria", [
        "Integration capability with existing ERP system",
        "User adoption and training requirements",
        "Total cost of ownership over 5 years",
        "Data migration complexity and risk",
        "Vendor support and SLA guarantees",
        "Customization flexibility for regional teams",
    ])

    # --- Slide 4: Stakeholder Input ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide4, "Stakeholder Input", [
        "Engineering: Prioritizes API extensibility and documentation quality",
        "Sales: Needs mobile access and offline capability",
        "Finance: Concerned about hidden licensing costs",
        "Operations: Requires automated workflow triggers",
        "HR: Wants seamless onboarding module integration",
    ])

    # --- Slide 5: Risk Assessment ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide5, "Risk Assessment", [
        "Vendor lock-in potential: Medium to High",
        "Data loss during migration: Low with proper ETL pipeline",
        "User resistance to change: Moderate, mitigated by training",
        "Timeline overrun risk: High if customization scope creeps",
        "Budget impact of delayed decision: ~$35K/month in legacy costs",
    ])

    # --- Slide 6: Pros and Cons Analysis (TITLE ONLY, no table) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Pros and Cons Analysis"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2F, 0x2F, 0x2F)

    # --- Slide 7: Next Steps & Timeline ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide7, "Next Steps & Timeline", [
        "April 7: Final vendor demos completed",
        "April 14: Scoring workshop with all stakeholders",
        "April 21: Executive review and selection",
        "May 1: Contract negotiation begins",
        "June 15: Implementation kickoff target",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
