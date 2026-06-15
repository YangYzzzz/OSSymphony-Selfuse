"""
Reward Script: Add a 2x4 table on slide 6 with colored bold headers
Task ID: impress_tct_031
Domain: libreoffice_impress
Scoring:
  - Component 1: Table exists on slide 6 with correct dimensions (0.25)
  - Component 2: Cell(0,0) text is 'Pros' (0.15)
  - Component 3: Cell(0,0) font is bold and green #2E7D32 (0.20)
  - Component 4: Cell(0,1) text is 'Cons' (0.15)
  - Component 5: Cell(0,1) font is bold and red #C62828 (0.20)
  - Component 6: Rows 2-4 are empty (0.05)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_031'


def persist_app_state(domain):
    """Best-effort save via Ctrl+S for unsaved GUI edits."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            import time
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def find_table_on_slide(slide):
    """Find the first table shape on a slide."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return shape.table
    return None


def get_cell_font_props(cell):
    """Get font properties from the first non-empty run in a cell."""
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            if (run.text or "").strip():
                bold = run.font.bold
                # Normalize None to False
                bold = False if bold is None else bold
                try:
                    rgb = str(run.font.color.rgb) if run.font.color.type is not None else None
                except Exception:
                    rgb = None
                return bold, rgb
    return None, None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide6 = prs.slides[5]  # 0-indexed

    # Component 1: Table exists on slide 6 with 4 rows x 2 columns (0.25 points)
    try:
        table = find_table_on_slide(slide6)
        if table is not None:
            rows = len(table.rows)
            cols = len(table.columns)
            if rows == 4 and cols == 2:
                print(f"PASS: Component 1 -- Table found on slide 6 with {rows}x{cols} dimensions (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 -- Table dimensions are {rows}x{cols}, expected 4x2")
        else:
            print("FAIL: Component 1 -- No table found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Stop if no table found (remaining checks depend on it)
    if table is None:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Cell(0,0) text is 'Pros' (0.15 points)
    try:
        cell_00_text = table.cell(0, 0).text.strip()
        if cell_00_text == 'Pros':
            print(f"PASS: Component 2 -- Cell(0,0) text is 'Pros' (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 -- Cell(0,0) text is '{cell_00_text}', expected 'Pros'")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Cell(0,0) font is bold and green #2E7D32 (0.20 points)
    try:
        bold_00, rgb_00 = get_cell_font_props(table.cell(0, 0))
        if bold_00 and rgb_00 == '2E7D32':
            print(f"PASS: Component 3 -- Cell(0,0) is bold with color #2E7D32 (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 3 -- Cell(0,0) bold={bold_00}, color={rgb_00}, expected bold=True, color=2E7D32")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Cell(0,1) text is 'Cons' (0.15 points)
    try:
        cell_01_text = table.cell(0, 1).text.strip()
        if cell_01_text == 'Cons':
            print(f"PASS: Component 4 -- Cell(0,1) text is 'Cons' (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 -- Cell(0,1) text is '{cell_01_text}', expected 'Cons'")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Cell(0,1) font is bold and red #C62828 (0.20 points)
    try:
        bold_01, rgb_01 = get_cell_font_props(table.cell(0, 1))
        if bold_01 and rgb_01 == 'C62828':
            print(f"PASS: Component 5 -- Cell(0,1) is bold with color #C62828 (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 5 -- Cell(0,1) bold={bold_01}, color={rgb_01}, expected bold=True, color=C62828")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # Component 6: Rows 2-4 are empty (0.05 points)
    try:
        all_empty = True
        for r in range(1, 4):
            for c in range(2):
                cell_text = table.cell(r, c).text.strip()
                if cell_text:
                    all_empty = False
                    print(f"  NOTE: cell({r},{c}) is not empty: '{cell_text}'")
        if all_empty:
            print(f"PASS: Component 6 -- Rows 2-4 are all empty (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 6 -- Some cells in rows 2-4 are not empty")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
