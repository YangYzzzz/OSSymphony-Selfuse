"""
Reward Script: Split merged cell in row 3 of table on slide 2
Task ID: impress_tct_015
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Cell (3,1) gridSpan is 1 (no horizontal merge)
  Component 2 (0.3): Cells (3,2) and (3,3) have no hMerge attribute
  Component 3 (0.3): Cell (3,1) retains the original merged-cell text
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_015'


def persist_app_state(domain: str):
    """Save any unsaved GUI edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify that the merged cell in row 3 of the table on slide 2
    has been split back into individual cells.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 2 slides
    if len(prs.slides) < 2:
        print(f"CRITICAL: Expected at least 2 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Find the table on slide 2 (index 1)
    slide = prs.slides[1]
    table = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            break

    if table is None:
        print("CRITICAL: No table found on slide 2")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: table has at least 4 rows and 4 columns
    if len(table.rows) < 4 or len(table.columns) < 4:
        print(f"CRITICAL: Table too small ({len(table.rows)}x{len(table.columns)}), expected at least 4x4")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Cell (3,1) gridSpan is 1 — no horizontal merge (0.4 points)
    # In initial state, gridSpan=3. After splitting, it should be 1 (or absent).
    try:
        cell_3_1 = table.cell(3, 1)
        tc_3_1 = cell_3_1._tc
        grid_span = tc_3_1.get('gridSpan', '1')
        if str(grid_span) == '1':
            print(f"PASS: Component 1 — Cell (3,1) gridSpan={grid_span}, not merged (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 — Cell (3,1) gridSpan={grid_span}, expected 1 (still merged)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Cells (3,2) and (3,3) have no hMerge attribute (0.3 points)
    # In initial state, these cells have hMerge=1 (they are continuation cells of the merge).
    # After splitting, hMerge should be absent or None.
    try:
        cell_3_2 = table.cell(3, 2)
        cell_3_3 = table.cell(3, 3)
        tc_3_2 = cell_3_2._tc
        tc_3_3 = cell_3_3._tc

        hmerge_2 = tc_3_2.get('hMerge', None)
        hmerge_3 = tc_3_3.get('hMerge', None)

        no_hmerge_2 = hmerge_2 is None or str(hmerge_2) == '0' or str(hmerge_2).lower() == 'false'
        no_hmerge_3 = hmerge_3 is None or str(hmerge_3) == '0' or str(hmerge_3).lower() == 'false'

        if no_hmerge_2 and no_hmerge_3:
            print(f"PASS: Component 2 — Cells (3,2) and (3,3) have no hMerge (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 — hMerge still present: (3,2)={hmerge_2}, (3,3)={hmerge_3}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Cell (3,1) retains original text AND is unmerged (0.3 points)
    # Compound check: text presence alone is a precondition (true in both envs).
    # We require BOTH that the text is present AND that the cell is no longer merged.
    try:
        cell_3_1_text = table.cell(3, 1).text.strip()
        expected_fragment = "Restructuring in progress"
        tc_check = table.cell(3, 1)._tc
        gs_check = str(tc_check.get('gridSpan', '1'))
        text_ok = expected_fragment in cell_3_1_text
        unmerged_ok = gs_check == '1'
        if text_ok and unmerged_ok:
            print(f"PASS: Component 3 — Cell (3,1) contains text AND is unmerged (0.3 pts)")
            total_score += 0.3
        elif not text_ok:
            print(f"FAIL: Component 3 — Cell (3,1) text missing: '{cell_3_1_text[:60]}'")
        else:
            print(f"FAIL: Component 3 — Cell (3,1) still merged (gridSpan={gs_check}), text present but merge not split")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
