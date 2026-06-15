"""
Initial Setup: Department Overview presentation with merged cell in table
Task ID: impress_tct_015
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_015'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Department Overview"
    slide1.placeholders[1].text = "FY 2025 Quarterly Review\nAcme Corporation"

    # ---- Slide 2: Table with merged cell in row 3 ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Add title textbox
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Staffing Summary by Department"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Table: 6 rows x 4 columns
    rows, cols = 6, 4
    tbl_shape = slide2.shapes.add_table(
        rows, cols,
        Inches(0.8), Inches(1.4), Inches(10.5), Inches(4.5)
    )
    table = tbl_shape.table

    # Set column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.5)
    table.columns[3].width = Inches(3.0)

    # Header row (row 0)
    headers = ["Department", "Headcount", "Budget (USD)", "Status"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.bold = True
                run.font.size = Pt(14)
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Data rows
    data = [
        # Row 1
        ["Engineering", "48", "$2,450,000", "On Track"],
        # Row 2
        ["Marketing", "23", "$1,120,000", "Under Review"],
        # Row 3 - col 0 has dept name, cols 1-3 will be merged with a note
        ["Sales", "", "", ""],
        # Row 4
        ["Human Resources", "12", "$680,000", "On Track"],
        # Row 5
        ["Finance", "15", "$890,000", "Expanding"],
    ]

    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = Pt(12)

    # Set row 3 col 0 text
    table.cell(3, 0).text = "Sales"
    for para in table.cell(3, 0).text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(12)

    # Merge cells in row 3: columns 1 through 3 (0-indexed: row index 3, cols 1-3)
    # python-pptx merge: cell(row, start_col).merge(cell(row, end_col))
    table.cell(3, 1).merge(table.cell(3, 3))
    # Set the merged cell text
    merged_cell = table.cell(3, 1)
    merged_cell.text = "Restructuring in progress - figures pending Q3 audit"
    for para in merged_cell.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(12)
            run.font.italic = True
            run.font.color.rgb = RGBColor(0x99, 0x33, 0x00)

    # ---- Slide 3: Team Highlights ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Engineering Team Highlights"
    run3 = p3.runs[0]
    run3.font.size = Pt(28)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    bullets_box = slide3.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(10), Inches(4.5))
    btf = bullets_box.text_frame
    btf.word_wrap = True
    items = [
        "Successfully launched v3.2 platform migration ahead of schedule",
        "Cloud infrastructure costs reduced by 18% through optimization",
        "New CI/CD pipeline deployed, reducing build times by 40%",
        "Onboarded 8 new engineers in Q1, retention rate at 96%",
    ]
    for i, item in enumerate(items):
        if i == 0:
            p = btf.paragraphs[0]
        else:
            p = btf.add_paragraph()
        p.text = item
        p.space_after = Pt(12)
        for run in p.runs:
            run.font.size = Pt(16)

    # ---- Slide 4: Budget Overview ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Budget Allocation FY2025"
    run4 = p4.runs[0]
    run4.font.size = Pt(28)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    content4 = slide4.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(10), Inches(4.5))
    ctf4 = content4.text_frame
    ctf4.word_wrap = True
    lines = [
        "Total Annual Budget: $12,400,000",
        "Engineering: $2,450,000 (19.8%)",
        "Marketing: $1,120,000 (9.0%)",
        "Sales: TBD pending restructuring",
        "Human Resources: $680,000 (5.5%)",
        "Finance: $890,000 (7.2%)",
        "Operations & Facilities: $3,260,000 (26.3%)",
        "Contingency Reserve: $4,000,000 (32.3%)",
    ]
    for i, line in enumerate(lines):
        if i == 0:
            p = ctf4.paragraphs[0]
        else:
            p = ctf4.add_paragraph()
        p.text = line
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.size = Pt(14)
        if i == 0:
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(18)

    # ---- Slide 5: Next Steps ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf5 = txBox5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Next Steps & Action Items"
    run5 = p5.runs[0]
    run5.font.size = Pt(28)
    run5.font.bold = True
    run5.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    content5 = slide5.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(10), Inches(4.5))
    ctf5 = content5.text_frame
    ctf5.word_wrap = True
    actions = [
        "Complete Sales department restructuring by end of Q3",
        "Finalize hiring plans for Engineering expansion",
        "Submit budget revision proposal to board by August 15",
        "Schedule mid-year performance reviews for all departments",
        "Review vendor contracts expiring in September",
    ]
    for i, action in enumerate(actions):
        if i == 0:
            p = ctf5.paragraphs[0]
        else:
            p = ctf5.add_paragraph()
        p.text = action
        p.space_after = Pt(12)
        for run in p.runs:
            run.font.size = Pt(16)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
