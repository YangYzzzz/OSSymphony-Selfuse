"""
Initial Setup: Create a 15-slide corporate training deck with white backgrounds and dark text.
Task ID: impress_gf5_004
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_004'
OUTPUT = f'{WORKDIR}/training_deck.pptx'

DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
SUBTITLE_TEXT = RGBColor(0x44, 0x44, 0x44)
WHITE_BG = RGBColor(0xFF, 0xFF, 0xFF)


def set_text_color(shape, color):
    """Set all text runs in a shape to the given color."""
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = color


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE_BG

    slide.shapes.title.text = title_text
    set_text_color(slide.shapes.title, DARK_TEXT)

    slide.placeholders[1].text = subtitle_text
    set_text_color(slide.placeholders[1], SUBTITLE_TEXT)
    return slide


def add_content_slide(prs, title_text, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE_BG

    slide.shapes.title.text = title_text
    set_text_color(slide.shapes.title, DARK_TEXT)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        for run in p.runs:
            run.font.color.rgb = DARK_TEXT
            run.font.size = Pt(18)
    return slide


def add_blank_text_slide(prs, title_text, body_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE_BG

    # Title text box
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = DARK_TEXT

    # Body text box
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8), Inches(4.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = body_text
    for run in p2.runs:
        run.font.size = Pt(16)
        run.font.color.rgb = DARK_TEXT
    return slide


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # Slide 1: Title slide
    add_title_slide(prs,
        "Corporate Training Program 2025",
        "Meridian Technologies Inc. - Learning & Development")

    # Slide 2: Agenda
    add_content_slide(prs, "Training Agenda", [
        "Company Overview & Mission Statement",
        "Department Structure and Key Contacts",
        "Compliance and Safety Protocols",
        "Technology Tools and Systems Access",
        "Performance Review Process",
    ])

    # Slide 3: Company Overview
    add_content_slide(prs, "Company Overview", [
        "Founded in 2008 with headquarters in Austin, TX",
        "Over 4,200 employees across 12 global offices",
        "Annual revenue of $1.8 billion (FY2024)",
        "Industry leader in enterprise cloud solutions",
    ])

    # Slide 4: Mission & Values
    add_content_slide(prs, "Mission & Values", [
        "Mission: Empower organizations through innovative technology",
        "Integrity: Transparent and ethical in every interaction",
        "Innovation: Constantly push the boundaries of what is possible",
        "Collaboration: Success is a team effort",
        "Customer Focus: Every decision starts with the customer",
    ])

    # Slide 5: Organizational Structure
    add_content_slide(prs, "Organizational Structure", [
        "CEO: Patricia Navarro",
        "VP Engineering: David Kim - 850 engineers",
        "VP Sales: Rachel Thornton - 620 representatives",
        "VP Marketing: James Okafor - 280 specialists",
        "VP HR: Lisa Fernandez - 95 HR professionals",
    ])

    # Slide 6: Key Departments
    add_content_slide(prs, "Key Departments & Contacts", [
        "IT Helpdesk: ext. 4500 or helpdesk@meridian.com",
        "Facilities Management: ext. 4200, Building A Room 102",
        "Human Resources: ext. 4100, 3rd Floor East Wing",
        "Security Office: ext. 4800, Main Lobby",
    ])

    # Slide 7: Compliance
    add_content_slide(prs, "Compliance Requirements", [
        "Annual data privacy training (GDPR & CCPA)",
        "Quarterly cybersecurity awareness modules",
        "Workplace harassment prevention certification",
        "Export control and trade compliance review",
        "Financial reporting ethics (SOX compliance)",
    ])

    # Slide 8: Safety Protocols
    add_content_slide(prs, "Safety Protocols", [
        "Emergency evacuation routes posted on every floor",
        "Fire drill scheduled first Monday of each quarter",
        "First aid stations located in break rooms on each floor",
        "Report safety hazards via the SafeTrack portal",
        "Incident response team available 24/7 at ext. 9911",
    ])

    # Slide 9: Technology Access
    add_content_slide(prs, "Technology Tools & Access", [
        "Laptop setup via IT onboarding portal (day 1)",
        "VPN access: GlobalConnect v4.2 - request through ServiceNow",
        "Collaboration: Microsoft Teams and Confluence",
        "Project Management: Jira and Asana",
        "Code repositories: GitHub Enterprise",
    ])

    # Slide 10: Email & Communication
    add_content_slide(prs, "Communication Guidelines", [
        "Corporate email via Outlook 365",
        "Response time SLA: internal 4 hours, external 24 hours",
        "Use Teams channels for project discussions",
        "All-hands meeting: last Thursday of each month",
        "Newsletter submissions due by the 15th",
    ])

    # Slide 11: Performance Reviews
    add_content_slide(prs, "Performance Review Process", [
        "Quarterly check-ins with direct manager",
        "Annual 360-degree review in November",
        "Self-assessment due two weeks before review",
        "Goals set using OKR framework each January",
        "Promotion cycles: March and September",
    ])

    # Slide 12: Benefits Overview
    add_content_slide(prs, "Benefits Overview", [
        "Health insurance: PPO and HMO plans available",
        "401(k) matching up to 6% of salary",
        "20 days PTO + 10 company holidays",
        "Tuition reimbursement: up to $8,000/year",
        "Wellness stipend: $1,200 annual",
    ])

    # Slide 13: Professional Development
    add_content_slide(prs, "Professional Development", [
        "LinkedIn Learning access for all employees",
        "Conference attendance budget: $3,500/year",
        "Internal mentorship program (apply in Q1)",
        "Tech talks every Wednesday at 3 PM",
        "Leadership development track for senior ICs",
    ])

    # Slide 14: FAQ
    add_content_slide(prs, "Frequently Asked Questions", [
        "Where do I park? Garage B, levels 2-4 for employees",
        "How do I request time off? Through the Workday portal",
        "Can I work remotely? Up to 3 days/week with manager approval",
        "Who do I contact for payroll issues? payroll@meridian.com",
    ])

    # Slide 15: Closing / Contact
    add_title_slide(prs,
        "Welcome to Meridian Technologies!",
        "Questions? Contact L&D at training@meridian.com | ext. 4350")

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
