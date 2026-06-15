"""
Reward Script: Change presentation theme to dark blue background with white text
Task ID: impress_gf5_004
Domain: libreoffice_impress
Scoring:
  Component 1: Slide count is 15 (0.1 pts)
  Component 2: All 15 slides have dark blue (#1B3A6B) background (0.45 pts, proportional)
  Component 3: All text runs across all slides are white (#FFFFFF) (0.45 pts, proportional)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_004'
EXPECTED_BG_RGB = '1B3A6B'
EXPECTED_TEXT_RGB = 'FFFFFF'
EXPECTED_SLIDE_COUNT = 15


def persist_app_state(domain: str):
    """Best-effort save any unsaved GUI edits."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)
    num_slides = len(slides)

    # Component 1: Slide count is 15 (0.1 points)
    # This is a precondition gate — if the file doesn't have 15 slides,
    # it's been corrupted. But we only award points if the count changed
    # (it should be 15 in both envs, so we use it as a gate, not scoring).
    # Actually: both initial and golden have 15 slides, so this would pass
    # on both. Instead, use slide count as a gate and redistribute points.
    # Re-design: only score bg and text changes (the actual task).

    # Component 1: Background color — all 15 slides have #1B3A6B (0.50 points)
    # Award proportionally: 0.50 / 15 per slide with correct bg
    try:
        bg_pass_count = 0
        for i, slide in enumerate(slides):
            try:
                fill = slide.background.fill
                if fill.type == 1:  # SOLID fill
                    bg_rgb = str(fill.fore_color.rgb)
                    if bg_rgb.upper() == EXPECTED_BG_RGB:
                        bg_pass_count += 1
                    else:
                        print(f"FAIL: Slide {i+1} bg color is {bg_rgb}, expected {EXPECTED_BG_RGB}")
                else:
                    print(f"FAIL: Slide {i+1} bg fill type is {fill.type}, expected SOLID (1)")
            except Exception as e:
                print(f"ERROR: Slide {i+1} bg check: {e}")

        if bg_pass_count > 0:
            bg_score = 0.50 * (bg_pass_count / EXPECTED_SLIDE_COUNT)
            total_score += bg_score
            print(f"PASS: Component 1 — {bg_pass_count}/{EXPECTED_SLIDE_COUNT} slides have correct bg #{EXPECTED_BG_RGB} ({bg_score:.3f} pts)")
        else:
            print(f"FAIL: Component 1 — 0/{EXPECTED_SLIDE_COUNT} slides have correct bg #{EXPECTED_BG_RGB}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Text color — all text runs on all slides are #FFFFFF (0.50 points)
    # Award proportionally per slide: for each slide, check that ALL non-empty text runs are white.
    # Score = 0.50 * (slides_with_all_white_text / 15)
    try:
        text_pass_count = 0
        for i, slide in enumerate(slides):
            slide_all_white = True
            has_text = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if not (run.text or "").strip():
                                continue
                            has_text = True
                            try:
                                if run.font.color.type is not None:
                                    rgb = str(run.font.color.rgb).upper()
                                    if rgb != EXPECTED_TEXT_RGB:
                                        slide_all_white = False
                                        print(f"FAIL: Slide {i+1} has text run with color {rgb}, expected {EXPECTED_TEXT_RGB}")
                                        break
                                else:
                                    # Inherited color — not explicitly set to white
                                    slide_all_white = False
                                    print(f"FAIL: Slide {i+1} has text run with inherited (None) color, expected explicit {EXPECTED_TEXT_RGB}")
                                    break
                            except Exception as e:
                                slide_all_white = False
                                print(f"ERROR: Slide {i+1} text color check: {e}")
                                break
                    if not slide_all_white:
                        break
            if has_text and slide_all_white:
                text_pass_count += 1
            elif not has_text:
                # Slide with no text — count as pass (nothing to change)
                text_pass_count += 1

        if text_pass_count > 0:
            text_score = 0.50 * (text_pass_count / EXPECTED_SLIDE_COUNT)
            total_score += text_score
            print(f"PASS: Component 2 — {text_pass_count}/{EXPECTED_SLIDE_COUNT} slides have all white text ({text_score:.3f} pts)")
        else:
            print(f"FAIL: Component 2 — 0/{EXPECTED_SLIDE_COUNT} slides have all white text")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score:.3f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/training_deck.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
