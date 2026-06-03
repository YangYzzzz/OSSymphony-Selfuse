"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 106 I’ve got three images that are all different widths and look off-kilter. In LibreOffice Impress, what’s the quickest way to line them up so they’re centered on the slide and each one is exactly 10 cm wide?
Generated: 2025-09-11 00:31:42
Status: success
Model: azure-o3
Total Steps: 15
"""

import os
from pptx import Presentation

def find_default_pptx():
    """Locate the task presentation in /home/user automatically (best-effort)."""
    search_root = '/home/user'
    for root, _, files in os.walk(search_root):
        for name in files:
            if name.endswith('.pptx') and 'on_slide_106' in name:
                return os.path.join(root, name)
    return None

def verify_align_and_resize(file_path: str) -> float:
    """Verify that on slide 106 there are exactly three images, each 10 cm wide and centred.

    Scoring (progressive):
        0.3 – Slide 106 contains exactly three picture shapes.
        0.3 – Every picture’s width is 10 cm ±1 %.
        0.4 – Every picture is horizontally centred on the slide (±20 000 EMU).
    Partial credit is granted proportionally for the width and centring checks.
    The function prints a detailed breakdown and returns a float between 0.0 and 1.0.
    """
    print(f'Verifying presentation: {file_path}')

    if not os.path.exists(file_path):
        print('✗ File not found')
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as exc:
        print(f'✗ Unable to load PPTX: {exc}')
        return 0.0

    # Slide 106 → 0-based index 105
    slide_index = 105
    if slide_index >= len(prs.slides):
        print(f'✗ Slide 106 is missing (presentation has {len(prs.slides)} slides)')
        return 0.0

    slide = prs.slides[slide_index]

    # Collect picture shapes (shape_type 13 corresponds to PICTURE)
    pictures = [sh for sh in slide.shapes if sh.shape_type == 13]
    print(f'Found {len(pictures)} picture shape(s) on slide 106')

    score = 0.0

    # 1) Correct number of images (0.3)
    if len(pictures) == 3:
        print('✓ Exactly three pictures found (+0.3)')
        score += 0.3
    elif len(pictures) > 0:
        partial = 0.3 * min(len(pictures), 3) / 3
        print(f'Partial credit for picture count (+{partial:.2f})')
        score += partial
    else:
        print('✗ No pictures – 0 points for this criterion')
        return 0.0  # nothing else to evaluate

    # 2) Width ≈ 10 cm (0.3)
    target_width = 3_600_000  # 10 cm in EMUs
    tolerance = target_width * 0.01  # ±1 %
    width_flags = [abs(sh.width - target_width) <= tolerance for sh in pictures]

    if all(width_flags):
        print('✓ All pictures are 10 cm wide ±1 % (+0.3)')
        score += 0.3
    else:
        correct = sum(width_flags)
        if correct:
            partial = 0.3 * correct / len(pictures)
            print(f'Partial width credit: {correct}/{len(pictures)} correct (+{partial:.2f})')
            score += partial
        else:
            print('✗ No picture has the correct width – 0 for this criterion')

    # 3) Horizontally centred (0.4)
    slide_center = prs.slide_width / 2
    center_tol = 20_000  # ≈0.055 cm
    center_flags = [abs((sh.left + sh.width / 2) - slide_center) <= center_tol for sh in pictures]

    if all(center_flags):
        print('✓ All pictures are horizontally centred (+0.4)')
        score += 0.4
    else:
        correct = sum(center_flags)
        if correct:
            partial = 0.4 * correct / len(pictures)
            print(f'Partial centring credit: {correct}/{len(pictures)} centred (+{partial:.2f})')
            score += partial
        else:
            print('✗ No picture is properly centred – 0 for this criterion')

    final_score = round(min(score, 1.0), 2)
    print(f'Total score: {final_score}')
    print(f'REWARD: {final_score}')
    return final_score

if __name__ == '__main__':
    import sys
    path = sys.argv[1] if len(sys.argv) > 1 else find_default_pptx()
    if path is None:
        print('✗ Presentation file not provided and could not be located automatically.')
        print('REWARD: 0.0')
    else:
        verify_align_and_resize(path)
