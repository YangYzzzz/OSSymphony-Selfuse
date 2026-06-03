"""
Reward Script: Paste unformatted text onto slide 4
Task ID: impress_tct_082
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Slide 4 has additional paragraph(s) compared to initial 5
  Component 2 (0.3): Pasted text contains expected content about remaining budget
  Component 3 (0.3): Pasted text uses destination formatting (Calibri ~14pt, not bold,
                      not red) — verifies unformatted paste, not source formatting
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_082'

# Number of paragraphs on slide 4 in the initial state (before paste)
INITIAL_SLIDE4_PARA_COUNT = 5

# Expected pasted text content (key phrase)
EXPECTED_PASTED_PHRASE = "influencer partnerships"

# Source formatting that should NOT appear if pasted as unformatted
SOURCE_FONT_NAME = "Times New Roman"
SOURCE_FONT_SIZE = 228600  # 18pt in EMU
SOURCE_FONT_BOLD = True

def persist_app_state(domain):
    """Send Ctrl+S to save any unsaved GUI edits."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide4 = prs.slides[3]  # 0-indexed

    # Find the content placeholder on slide 4
    content_shape = None
    for shape in slide4.shapes:
        if hasattr(shape, 'text_frame') and shape.name != 'Title 1':
            content_shape = shape
            break

    if content_shape is None:
        print("FAIL: No content shape found on slide 4")
        print("REWARD: 0.0")
        return 0.0

    tf = content_shape.text_frame
    para_count = len(tf.paragraphs)
    print(f"INFO: Slide 4 content has {para_count} paragraphs (initial had {INITIAL_SLIDE4_PARA_COUNT})")

    # Component 1: Slide 4 has additional paragraph(s) beyond the initial 5 (0.4 points)
    # This verifies that new text was actually pasted onto slide 4
    try:
        if para_count > INITIAL_SLIDE4_PARA_COUNT:
            print(f"PASS: Component 1 — Slide 4 has {para_count} paragraphs, more than initial {INITIAL_SLIDE4_PARA_COUNT} (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 — Slide 4 still has {para_count} paragraphs, expected more than {INITIAL_SLIDE4_PARA_COUNT}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Identify the new paragraphs (those beyond the initial count)
    new_paragraphs = tf.paragraphs[INITIAL_SLIDE4_PARA_COUNT:] if para_count > INITIAL_SLIDE4_PARA_COUNT else []
    new_text = " ".join(p.text for p in new_paragraphs).strip()
    print(f"INFO: New paragraph text: [{new_text[:100]}]")

    # Component 2: Pasted text contains expected content (0.3 points)
    # The clipboard had text about remaining budget allocation
    try:
        if new_text and EXPECTED_PASTED_PHRASE.lower() in new_text.lower():
            print(f"PASS: Component 2 — Pasted text contains '{EXPECTED_PASTED_PHRASE}' (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 — Expected '{EXPECTED_PASTED_PHRASE}' in new text, found: [{new_text[:80]}]")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Pasted text uses destination formatting, NOT source formatting (0.3 points)
    # Source was: 18pt Times New Roman, bold, red
    # Destination is: 14pt Calibri, not bold, no explicit color
    # If pasted as unformatted text, it should match the destination style
    try:
        if not new_paragraphs:
            print("FAIL: Component 3 — No new paragraphs to check formatting")
        else:
            formatting_ok = True
            issues = []

            for p_idx, para in enumerate(new_paragraphs):
                for run in para.runs:
                    if not (run.text or "").strip():
                        continue

                    # Check font name: should be Calibri (destination), not Times New Roman (source)
                    fname = run.font.name
                    if fname == SOURCE_FONT_NAME:
                        formatting_ok = False
                        issues.append(f"Run has source font '{SOURCE_FONT_NAME}' instead of destination font")

                    # Check font size: should be ~14pt (177800 EMU), not 18pt (228600 EMU)
                    fsize = run.font.size
                    if fsize is not None and fsize == SOURCE_FONT_SIZE:
                        formatting_ok = False
                        issues.append(f"Run has source size 18pt instead of destination ~14pt")

                    # Check bold: should not be bold (source was bold)
                    fbold = run.font.bold
                    if fbold is True:
                        formatting_ok = False
                        issues.append(f"Run is bold (source formatting leaked)")

                    # Check color: should not be red (source was red)
                    try:
                        if run.font.color.type is not None:
                            rgb = str(run.font.color.rgb)
                            if rgb.upper() == "FF0000":
                                formatting_ok = False
                                issues.append(f"Run has red color (source formatting leaked)")
                    except Exception:
                        pass  # No color set is fine

            if formatting_ok:
                print(f"PASS: Component 3 — Pasted text uses destination formatting (not source 18pt/TNR/bold/red) (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 3 — Source formatting detected: {'; '.join(issues)}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
