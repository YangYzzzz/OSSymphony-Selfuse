"""
Initial Setup: Paste unformatted text onto slide 4
Task ID: impress_tct_082
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_082'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q3 Marketing Strategy Review"
    slide1.placeholders[1].text = "Prepared by Elena Vasquez\nSeptember 2025"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Campaign Performance Overview"
    for item in ["Budget Allocation Updates", "Customer Acquisition Metrics",
                  "Social Media Engagement Results", "Next Quarter Planning"]:
        p = tf2.add_paragraph()
        p.text = item
        p.level = 0
    for para in tf2.paragraphs:
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(14)

    # --- Slide 3: Campaign Performance ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Campaign Performance"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Email campaigns achieved 24.3% open rate"
    for item in ["Social media impressions up 18% quarter-over-quarter",
                  "Paid search ROI improved to 3.2x",
                  "Content marketing generated 1,247 qualified leads"]:
        p = tf3.add_paragraph()
        p.text = item
        p.level = 0
    for para in tf3.paragraphs:
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(14)

    # --- Slide 4: Budget Allocation (the target slide) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Budget Allocation"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Current spending breakdown for Q3:"
    items4 = [
        "Digital Advertising: $45,200 (32%)",
        "Content Production: $28,750 (20%)",
        "Events and Sponsorships: $21,300 (15%)",
        "Email Platform and Tools: $12,800 (9%)",
    ]
    for item in items4:
        p = tf4.add_paragraph()
        p.text = item
        p.level = 0
    # Set all existing text to 14pt Calibri
    for para in tf4.paragraphs:
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(14)

    # --- Slide 5: Next Steps ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Next Steps"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Finalize Q4 budget proposal by October 15"
    for item in ["Schedule stakeholder review meeting",
                  "Complete competitive analysis report",
                  "Launch new social media pilot program"]:
        p = tf5.add_paragraph()
        p.text = item
        p.level = 0
    for para in tf5.paragraphs:
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(14)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Set clipboard with formatted text (18pt Times New Roman, bold, red)
    # We use xclip to set plain text on clipboard; the task context says
    # clipboard has formatted text from Word doc, but on Linux VM we set it as plain text.
    clipboard_text = "Remaining 24% allocated to influencer partnerships and affiliate marketing programs across three regional markets"
    try:
        proc = subprocess.Popen(
            ['xclip', '-selection', 'clipboard'],
            stdin=subprocess.PIPE,
            env={**os.environ, 'DISPLAY': ':0'}
        )
        proc.communicate(input=clipboard_text.encode('utf-8'))
        print(f'Clipboard set with text: {clipboard_text[:50]}...')
    except Exception as e:
        print(f'Warning: Could not set clipboard via xclip: {e}')
        # Fallback: use xsel
        try:
            proc = subprocess.Popen(
                ['xsel', '--clipboard', '--input'],
                stdin=subprocess.PIPE,
                env={**os.environ, 'DISPLAY': ':0'}
            )
            proc.communicate(input=clipboard_text.encode('utf-8'))
            print(f'Clipboard set via xsel')
        except Exception as e2:
            print(f'Warning: Could not set clipboard via xsel either: {e2}')

    # Also write the clipboard text to a temp file for reference
    with open('/tmp/clipboard_text.txt', 'w') as f:
        f.write(clipboard_text)

    # GUI-ready: open the presentation in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
