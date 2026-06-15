"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 125 clashes with the rest of the deck—can you switch its background to a solid fill, hex #FFD966 (the “Light Orange 1” swatch in Impress)?
Generated: 2025-09-10 23:34:18
Status: success
Model: azure-o3
Total Steps: 1
"""

from pptx import Presentation
from pptx.enum.dml import MSO_FILL
import os


def verify_slide_125_background(file_path: str) -> float:
    """
    Verify that slide 125 of the given PPTX has a **solid** background
    filled with the hex colour **#FFD966** (“Light Orange 1”).

    Progressive scoring (max 1.0):
      • 0.2 – File loads and has ≥ 125 slides (pre-requisite for the task)
      • 0.3 – Slide 125 background is a *solid* fill (not gradient/picture/etc.)
      • 0.5 – Solid fill colour exactly matches #FFD966

    Returns a float in [0.0, 1.0] and prints a detailed breakdown.
    """

    score = 0.0
    max_score = 1.0

    # ------------------------------------------------------------------
    # 1. File existence & load (no points – prerequisite)
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Failed to load PPTX: {e}")
        print("REWARD: 0.0")
        return 0.0

    # ------------------------------------------------------------------
    # 2. Ensure at least 125 slides exist (0.2 pts)
    # ------------------------------------------------------------------
    slide_count = len(prs.slides)
    print(f"Slide count detected: {slide_count}")
    if slide_count >= 125:
        score += 0.2
        print("✓ Presentation contains slide 125 (0.2)")
    else:
        print("✗ Presentation has fewer than 125 slides – cannot verify task")
        print(f"REWARD: {score}")
        return score  # cannot continue further

    # ------------------------------------------------------------------
    # 3. Inspect slide 125 background
    # ------------------------------------------------------------------
    target_slide = prs.slides[124]  # zero-based index
    fill = target_slide.background.fill

    # 3a. Check for SOLID fill type (0.3 pts)
    if fill.type == MSO_FILL.SOLID:
        score += 0.3
        print("✓ Slide 125 background is SOLID fill (0.3)")
    else:
        print("✗ Slide 125 background is not a solid fill – task incomplete")
        print(f"REWARD: {score}")
        return score

    # 3b. Verify hex colour #FFD966 (0.5 pts)
    rgb = getattr(fill.fore_color, 'rgb', None)
    if rgb is not None:
        rgb_hex = str(rgb).upper()
        print(f"Detected colour on slide 125: #{rgb_hex}")
        if rgb_hex == "FFD966":
            score += 0.5
            print("✓ Background colour matches #FFD966 (0.5)")
        else:
            print("✗ Background colour does NOT match #FFD966")
    else:
        print("✗ Could not determine RGB colour of the background")

    # ------------------------------------------------------------------
    # Final score (capped at 1.0)
    # ------------------------------------------------------------------
    final_score = min(score, max_score)
    print(f"Total score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score


# ----------------------------------------------------------------------
# Execute verification when script is run directly
# ----------------------------------------------------------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/slide_125_clashes_with_the_rest_of_the_deckcan_you_switch_its_background_to_a_solid_fill_hex_ffd966__golden.pptx"
    verify_slide_125_background(FILE_PATH)
