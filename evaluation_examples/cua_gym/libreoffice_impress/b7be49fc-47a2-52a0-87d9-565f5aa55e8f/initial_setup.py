"""
Initial Setup: Create a UX case study presentation with 5 content slides and a blank slide 6.
Also create before.jpg and after.jpg placeholder images for the comparison task.
Task ID: impress_gf5_039
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_039'
OUTPUT = f'{WORKDIR}/case_study.pptx'
IMG_DIR = f'{WORKDIR}/images'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_sample_images():
    """Create realistic before/after mockup images for the UX case study."""
    os.makedirs(IMG_DIR, exist_ok=True)

    # before.jpg - old UI design (muted, cluttered look)
    img_before = Image.new('RGB', (800, 600), color=(235, 235, 230))
    draw = ImageDraw.Draw(img_before)
    # Header bar
    draw.rectangle([0, 0, 800, 60], fill=(102, 102, 102))
    draw.text((20, 18), "Old Dashboard - v2.1", fill=(255, 255, 255))
    # Cluttered sidebar
    draw.rectangle([0, 60, 180, 600], fill=(220, 220, 215))
    for i in range(8):
        y = 80 + i * 55
        draw.rectangle([10, y, 170, y + 40], fill=(190, 190, 185))
        draw.text((20, y + 10), f"Menu Item {i+1}", fill=(80, 80, 80))
    # Content area with small cards
    for row in range(3):
        for col in range(3):
            x = 200 + col * 195
            y = 80 + row * 165
            draw.rectangle([x, y, x + 180, y + 150], fill=(245, 245, 240), outline=(200, 200, 200))
            draw.text((x + 10, y + 10), f"Widget {row*3+col+1}", fill=(100, 100, 100))
    img_before.save(f'{IMG_DIR}/before.jpg', 'JPEG', quality=90)

    # after.jpg - redesigned UI (clean, modern look)
    img_after = Image.new('RGB', (800, 600), color=(250, 251, 252))
    draw = ImageDraw.Draw(img_after)
    # Modern header
    draw.rectangle([0, 0, 800, 56], fill=(37, 99, 235))
    draw.text((24, 16), "Dashboard Pro", fill=(255, 255, 255))
    # Clean sidebar
    draw.rectangle([0, 56, 64, 600], fill=(30, 41, 59))
    for i in range(5):
        y = 76 + i * 50
        draw.rounded_rectangle([12, y, 52, y + 36], radius=6, fill=(51, 65, 85))
    # Modern cards
    for row in range(2):
        for col in range(2):
            x = 88 + col * 348
            y = 80 + row * 250
            draw.rounded_rectangle([x, y, x + 330, y + 230], radius=12, fill=(255, 255, 255), outline=(226, 232, 240))
            draw.text((x + 20, y + 20), ["Revenue", "Users", "Engagement", "Growth"][row*2+col], fill=(30, 41, 59))
            # Fake chart line
            draw.line([(x+20, y+180), (x+100, y+130), (x+200, y+150), (x+310, y+90)], fill=(37, 99, 235), width=3)
    img_after.save(f'{IMG_DIR}/after.jpg', 'JPEG', quality=90)

    print(f'Created {IMG_DIR}/before.jpg and {IMG_DIR}/after.jpg')


def add_text(slide, left, top, width, height, text, font_size=18, bold=False,
             color=RGBColor(0x33, 0x33, 0x33), alignment=PP_ALIGN.LEFT):
    """Helper to add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x2A, 0x4A)
    add_text(slide1, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
             "UX Redesign Case Study", font_size=40, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_text(slide1, Inches(1.5), Inches(3.8), Inches(10), Inches(1.0),
             "Dashboard Pro — From Legacy to Modern Experience", font_size=22,
             color=RGBColor(0xA0, 0xB4, 0xD0), alignment=PP_ALIGN.CENTER)
    add_text(slide1, Inches(1.5), Inches(5.2), Inches(10), Inches(0.6),
             "Elena Rodriguez, Senior UX Researcher  |  March 2025", font_size=16,
             color=RGBColor(0x80, 0x90, 0xA8), alignment=PP_ALIGN.CENTER)

    # --- Slide 2: Problem Statement ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide2, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "The Problem", font_size=32, bold=True,
             color=RGBColor(0x1B, 0x2A, 0x4A))
    problems = [
        "User satisfaction scores dropped 23% over the past two quarters",
        "Average task completion time increased from 45s to 72s",
        "Support ticket volume for navigation issues rose by 40%",
        "New user onboarding abandonment rate reached 35%",
    ]
    for i, prob in enumerate(problems):
        add_text(slide2, Inches(1.2), Inches(1.6 + i * 1.2), Inches(10), Inches(1.0),
                 f"• {prob}", font_size=18, color=RGBColor(0x44, 0x44, 0x44))

    # --- Slide 3: Research Methodology ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide3, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Research Methodology", font_size=32, bold=True,
             color=RGBColor(0x1B, 0x2A, 0x4A))
    methods = [
        ("User Interviews", "Conducted 24 in-depth interviews across 4 user segments"),
        ("Heuristic Evaluation", "3 UX experts assessed 48 interface touchpoints"),
        ("A/B Testing", "Ran 6 prototype variants with 1,200 participants over 3 weeks"),
        ("Analytics Review", "Analyzed 90 days of click-stream and session data"),
    ]
    for i, (title, desc) in enumerate(methods):
        add_text(slide3, Inches(1.0), Inches(1.5 + i * 1.3), Inches(4), Inches(0.5),
                 title, font_size=20, bold=True, color=RGBColor(0x25, 0x63, 0xEB))
        add_text(slide3, Inches(1.0), Inches(2.0 + i * 1.3), Inches(10), Inches(0.5),
                 desc, font_size=16, color=RGBColor(0x55, 0x55, 0x55))

    # --- Slide 4: Key Findings ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide4, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Key Findings", font_size=32, bold=True,
             color=RGBColor(0x1B, 0x2A, 0x4A))
    findings = [
        "78% of users found the navigation hierarchy confusing",
        "Card-based layouts improved scan-ability scores by 62%",
        "Reducing sidebar items from 8 to 5 cut task completion time by 38%",
        "Users preferred the blue accent palette (NPS +41 vs. control)",
        "Icon-only sidebar scored highest on space efficiency metrics",
    ]
    for i, finding in enumerate(findings):
        add_text(slide4, Inches(1.2), Inches(1.5 + i * 1.0), Inches(10), Inches(0.8),
                 f"{i+1}. {finding}", font_size=17, color=RGBColor(0x44, 0x44, 0x44))

    # --- Slide 5: Recommendations ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide5, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Recommendations", font_size=32, bold=True,
             color=RGBColor(0x1B, 0x2A, 0x4A))
    recs = [
        "Adopt a card-based dashboard layout with 4 primary widgets",
        "Implement a collapsible icon sidebar with 5 core navigation items",
        "Migrate to the blue accent color system (#2563EB primary)",
        "Add contextual tooltips for first-time user guidance",
    ]
    for i, rec in enumerate(recs):
        add_text(slide5, Inches(1.2), Inches(1.6 + i * 1.2), Inches(10), Inches(1.0),
                 f"→ {rec}", font_size=18, color=RGBColor(0x44, 0x44, 0x44))

    # --- Slide 6: Blank (for Before/After comparison) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # This slide is intentionally left blank for the agent to build the comparison

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_sample_images()
create_initial()
