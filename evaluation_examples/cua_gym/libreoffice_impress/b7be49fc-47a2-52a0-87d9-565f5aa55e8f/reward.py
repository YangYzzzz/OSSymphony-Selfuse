"""
Reward Script: Before/After comparison slide on slide 6 of case_study.pptx
Task ID: impress_gf5_039
Domain: libreoffice_impress
Scoring:
  Component 1 (0.15): BEFORE label on left half with bold text
  Component 2 (0.15): AFTER label on right half with bold text
  Component 3 (0.20): before.jpg image placed on left half
  Component 4 (0.20): after.jpg image placed on right half
  Component 5 (0.15): Vertical dividing line near center of slide
  Component 6 (0.15): Animation/timing exists on slide 6 (reveal on click)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_039'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, expected at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide6 = prs.slides[5]
    slide_width = prs.slide_width
    slide_center_x = slide_width // 2

    # Collect all shapes info
    text_shapes = []
    picture_shapes = []
    line_shapes = []

    for shape in slide6.shapes:
        if shape.has_text_frame:
            full_text = ""
            bold_runs = []
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    full_text += r.text
                    bold_runs.append(r.font.bold is True)
            text_shapes.append({
                'name': shape.name,
                'text': full_text.strip(),
                'bold': any(bold_runs),
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'center_x': shape.left + shape.width // 2,
            })
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture_shapes.append({
                'name': shape.name,
                'blob_len': len(shape.image.blob),
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'center_x': shape.left + shape.width // 2,
            })
        # LINE type is 9
        if shape.shape_type == 9:
            line_shapes.append({
                'name': shape.name,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
            })

    # Read reference image sizes for matching
    before_size = None
    after_size = None
    try:
        before_size = os.path.getsize('/home/user/images/before.jpg')
        after_size = os.path.getsize('/home/user/images/after.jpg')
    except Exception as e:
        print(f"WARN: Could not read reference image sizes: {e}")

    # Component 1: BEFORE label on left half with bold text (0.15 points)
    try:
        before_label = None
        for ts in text_shapes:
            if 'BEFORE' in ts['text'].upper():
                before_label = ts
                break
        if before_label is not None:
            # Check it's on the left half (center_x < slide center)
            on_left = before_label['center_x'] < slide_center_x
            if on_left and before_label['bold']:
                print(f"PASS: Component 1 - 'BEFORE' label found on left half, bold=True (0.15 pts)")
                total_score += 0.15
            elif on_left:
                print(f"PARTIAL: Component 1 - 'BEFORE' label on left but not bold (0.075 pts)")
                total_score += 0.075
            elif before_label['bold']:
                print(f"PARTIAL: Component 1 - 'BEFORE' label bold but not on left half (0.075 pts)")
                total_score += 0.075
            else:
                print(f"FAIL: Component 1 - 'BEFORE' label found but not bold and not on left half")
        else:
            print(f"FAIL: Component 1 - No 'BEFORE' label found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: AFTER label on right half with bold text (0.15 points)
    try:
        after_label = None
        for ts in text_shapes:
            if 'AFTER' in ts['text'].upper() and 'BEFORE' not in ts['text'].upper():
                after_label = ts
                break
        if after_label is not None:
            # Check it's on the right half (center_x > slide center)
            on_right = after_label['center_x'] > slide_center_x
            if on_right and after_label['bold']:
                print(f"PASS: Component 2 - 'AFTER' label found on right half, bold=True (0.15 pts)")
                total_score += 0.15
            elif on_right:
                print(f"PARTIAL: Component 2 - 'AFTER' label on right but not bold (0.075 pts)")
                total_score += 0.075
            elif after_label['bold']:
                print(f"PARTIAL: Component 2 - 'AFTER' label bold but not on right half (0.075 pts)")
                total_score += 0.075
            else:
                print(f"FAIL: Component 2 - 'AFTER' label found but not bold and not on right half")
        else:
            print(f"FAIL: Component 2 - No 'AFTER' label found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: before.jpg image placed on left half (0.20 points)
    try:
        # Match by blob size to before.jpg reference
        before_img = None
        for ps in picture_shapes:
            if before_size is not None and ps['blob_len'] == before_size:
                before_img = ps
                break
        if before_img is not None:
            on_left = before_img['center_x'] < slide_center_x
            if on_left:
                print(f"PASS: Component 3 - before.jpg image on left half (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 3 - before.jpg image found but NOT on left half (center_x={before_img['center_x']}, slide_center={slide_center_x})")
        else:
            # Fallback: if we have a picture on left half, give partial
            left_pics = [p for p in picture_shapes if p['center_x'] < slide_center_x]
            if left_pics:
                print(f"PARTIAL: Component 3 - A picture on left half but could not match before.jpg by size (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 3 - No picture found on left half of slide 6")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: after.jpg image placed on right half (0.20 points)
    try:
        after_img = None
        for ps in picture_shapes:
            if after_size is not None and ps['blob_len'] == after_size:
                after_img = ps
                break
        if after_img is not None:
            on_right = after_img['center_x'] > slide_center_x
            if on_right:
                print(f"PASS: Component 4 - after.jpg image on right half (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 4 - after.jpg image found but NOT on right half (center_x={after_img['center_x']}, slide_center={slide_center_x})")
        else:
            right_pics = [p for p in picture_shapes if p['center_x'] > slide_center_x]
            if right_pics:
                print(f"PARTIAL: Component 4 - A picture on right half but could not match after.jpg by size (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 4 - No picture found on right half of slide 6")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Vertical dividing line near center of slide (0.15 points)
    try:
        vertical_center_line = None
        for ls in line_shapes:
            # A vertical line has width ~0 and significant height
            is_vertical = ls['width'] <= 50000 and ls['height'] > 1000000  # essentially zero width, tall
            # Near center: left position close to slide center (within 15% of slide width)
            tolerance = slide_width * 0.15
            near_center = abs(ls['left'] - slide_center_x) < tolerance
            if is_vertical and near_center:
                vertical_center_line = ls
                break

        if vertical_center_line is not None:
            print(f"PASS: Component 5 - Vertical dividing line near center (left={vertical_center_line['left']}, center={slide_center_x}) (0.15 pts)")
            total_score += 0.15
        else:
            # Check if any line exists at all on slide 6
            if line_shapes:
                print(f"PARTIAL: Component 5 - Line(s) found but none vertical near center (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 5 - No line shape found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    # Component 6: Animation/timing exists on slide 6 with entrance effects (0.15 points)
    try:
        animation_found = 0  # 0=none, 1=timing only, 2=timing+entrance
        animated_shape_ids = set()

        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide6.xml') as f:
                root = ET.parse(f).getroot()
                ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'

                # Look for timing element
                timing = root.find(f'{{{ns_p}}}timing')
                if timing is not None:
                    animation_found = 1

                    # Look for entrance animations (presetClass="entr")
                    for ctn in root.iter(f'{{{ns_p}}}cTn'):
                        if ctn.get('presetClass') == 'entr':
                            animation_found = 2

                    # Get animated shape IDs
                    for spTgt in root.iter(f'{{{ns_p}}}spTgt'):
                        spid = spTgt.get('spid')
                        if spid:
                            animated_shape_ids.add(spid)

        if animation_found >= 2 and len(animated_shape_ids) >= 2:
            print(f"PASS: Component 6 - Animation with entrance effects on {len(animated_shape_ids)} shapes (0.15 pts)")
            total_score += 0.15
        elif animation_found >= 2:
            print(f"PARTIAL: Component 6 - Animation with entrance effect but only {len(animated_shape_ids)} animated shapes (0.10 pts)")
            total_score += 0.10
        elif animation_found >= 1:
            print(f"PARTIAL: Component 6 - Timing section exists but no entrance effects found (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 6 - No animation/timing found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 6 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Main entry point
file_path = f'{WORKDIR}/case_study.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
