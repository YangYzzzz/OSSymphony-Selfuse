"""
Reward Script: Apply bold + #1B3A6B color + underline to slide titles 3, 4, and 6
Task ID: osworld_impress_title_selective_formatting_014
Domain: libreoffice_impress

Scoring rubric:
  Component 1: Target slides 3, 4, 6 have bold=True AND untouched slides do NOT have
               newly added bold (selective bold applied correctly)              — 0.35 pts
  Component 2: Target slides 3, 4, 6 have color=#1B3A6B AND untouched slides
               do NOT have #1B3A6B color (selective color applied correctly)   — 0.40 pts
  Component 3: Target slides 3, 4, 6 have underline=True AND untouched slides
               do NOT have newly added underline (selective underline applied)  — 0.25 pts
  Total: 1.0

Notes:
  - Components fail on initial_env (all titles are plain/black, target slides lack formatting)
  - Components pass on golden_env (target slides formatted, untouched remain plain)
"""

import os

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_selective_formatting_014'

# Target slides that MUST be formatted (1-indexed)
TARGET_SLIDES = [3, 4, 6]
# Slides that MUST remain untouched (1-indexed)
UNTOUCHED_SLIDES = [1, 2, 5, 7]
# Required color
REQUIRED_COLOR = '1B3A6B'


def get_title_shape(slide):
    """Return the title placeholder (idx=0) for the given slide, or None."""
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.idx == 0:
                return shape
    return None


def get_nonempty_runs(title_shape):
    """Return a flat list of non-empty runs from the title text frame."""
    runs = []
    for para in title_shape.text_frame.paragraphs:
        for r in para.runs:
            if (r.text or '').strip():
                runs.append(r)
    return runs


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must be a 7-slide presentation
    if len(prs.slides) != 7:
        print(f"CRITICAL: Expected 7 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)

    # ------------------------------------------------------------------
    # Component 1: Selective bold — target slides bold=True, untouched slides
    # do not have bold=True  (0.35 pts)
    # Fails on initial_env because target slides lack bold.
    # ------------------------------------------------------------------
    try:
        comp1_failures = []

        # Target slides must have bold=True
        for slide_num in TARGET_SLIDES:
            slide = slides[slide_num - 1]
            title_shape = get_title_shape(slide)
            if title_shape is None:
                comp1_failures.append(f'Slide {slide_num}: no title shape')
                continue
            runs = get_nonempty_runs(title_shape)
            if not runs:
                comp1_failures.append(f'Slide {slide_num}: no text runs')
                continue
            all_bold = all(r.font.bold is True for r in runs)
            if not all_bold:
                comp1_failures.append(f'Slide {slide_num}: bold={[r.font.bold for r in runs]} (expected True)')
            else:
                print(f"  CHECK: Slide {slide_num}: bold=True OK")

        # Untouched slides must NOT have bold=True (guard against over-formatting)
        for slide_num in UNTOUCHED_SLIDES:
            slide = slides[slide_num - 1]
            title_shape = get_title_shape(slide)
            if title_shape is None:
                continue
            runs = get_nonempty_runs(title_shape)
            for r in runs:
                if r.font.bold is True:
                    comp1_failures.append(f'Slide {slide_num}: bold=True on untouched slide (should not be bold)')

        if not comp1_failures:
            print(f"PASS: Component 1 — selective bold applied correctly to slides 3, 4, 6 (0.35 pts)")
            total_score += 0.35
        else:
            for d in comp1_failures:
                print(f"  INFO: {d}")
            print(f"FAIL: Component 1 — selective bold check failed")
    except Exception as e:
        print(f"ERROR: Component 1 (bold check) — {e}")

    # ------------------------------------------------------------------
    # Component 2: Selective color — target slides color=#1B3A6B, untouched slides
    # do NOT have #1B3A6B color  (0.40 pts)
    # Fails on initial_env because target slides still have black (000000).
    # ------------------------------------------------------------------
    try:
        comp2_failures = []

        # Target slides must have color=1B3A6B
        for slide_num in TARGET_SLIDES:
            slide = slides[slide_num - 1]
            title_shape = get_title_shape(slide)
            if title_shape is None:
                comp2_failures.append(f'Slide {slide_num}: no title shape')
                continue
            runs = get_nonempty_runs(title_shape)
            if not runs:
                comp2_failures.append(f'Slide {slide_num}: no text runs')
                continue
            for r in runs:
                try:
                    if r.font.color.type is None:
                        comp2_failures.append(f'Slide {slide_num}: color=inherited/None (expected {REQUIRED_COLOR})')
                    else:
                        actual_color = str(r.font.color.rgb).upper()
                        if actual_color != REQUIRED_COLOR:
                            comp2_failures.append(
                                f'Slide {slide_num}: color={actual_color} (expected {REQUIRED_COLOR})')
                        else:
                            print(f"  CHECK: Slide {slide_num}: color={actual_color} OK")
                except Exception as ce:
                    comp2_failures.append(f'Slide {slide_num}: color error: {ce}')

        # Untouched slides must NOT have navy color
        for slide_num in UNTOUCHED_SLIDES:
            slide = slides[slide_num - 1]
            title_shape = get_title_shape(slide)
            if title_shape is None:
                continue
            runs = get_nonempty_runs(title_shape)
            for r in runs:
                try:
                    if r.font.color.type is not None:
                        actual_color = str(r.font.color.rgb).upper()
                        if actual_color == REQUIRED_COLOR:
                            comp2_failures.append(
                                f'Slide {slide_num}: color={actual_color} on untouched slide (should not be navy)')
                except Exception:
                    pass

        if not comp2_failures:
            print(f"PASS: Component 2 — selective color #{REQUIRED_COLOR} applied to slides 3, 4, 6 only (0.40 pts)")
            total_score += 0.40
        else:
            for d in comp2_failures:
                print(f"  INFO: {d}")
            print(f"FAIL: Component 2 — selective color check failed")
    except Exception as e:
        print(f"ERROR: Component 2 (color check) — {e}")

    # ------------------------------------------------------------------
    # Component 3: Selective underline — target slides underline=True, untouched
    # slides do NOT have underline=True  (0.25 pts)
    # Fails on initial_env because target slides lack underline.
    # ------------------------------------------------------------------
    try:
        comp3_failures = []

        # Target slides must have underline=True
        for slide_num in TARGET_SLIDES:
            slide = slides[slide_num - 1]
            title_shape = get_title_shape(slide)
            if title_shape is None:
                comp3_failures.append(f'Slide {slide_num}: no title shape')
                continue
            runs = get_nonempty_runs(title_shape)
            if not runs:
                comp3_failures.append(f'Slide {slide_num}: no text runs')
                continue
            all_underline = all(r.font.underline is True for r in runs)
            if not all_underline:
                comp3_failures.append(
                    f'Slide {slide_num}: underline={[r.font.underline for r in runs]} (expected True)')
            else:
                print(f"  CHECK: Slide {slide_num}: underline=True OK")

        # Untouched slides must NOT have underline=True
        for slide_num in UNTOUCHED_SLIDES:
            slide = slides[slide_num - 1]
            title_shape = get_title_shape(slide)
            if title_shape is None:
                continue
            runs = get_nonempty_runs(title_shape)
            for r in runs:
                if r.font.underline is True:
                    comp3_failures.append(
                        f'Slide {slide_num}: underline=True on untouched slide (should not be underlined)')

        if not comp3_failures:
            print(f"PASS: Component 3 — selective underline applied correctly to slides 3, 4, 6 (0.25 pts)")
            total_score += 0.25
        else:
            for d in comp3_failures:
                print(f"  INFO: {d}")
            print(f"FAIL: Component 3 — selective underline check failed")
    except Exception as e:
        print(f"ERROR: Component 3 (underline check) — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
