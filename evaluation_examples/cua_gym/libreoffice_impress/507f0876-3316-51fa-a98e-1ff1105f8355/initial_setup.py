"""
Initial Setup: 7-slide annual report deck with plain-formatted titles
Task ID: osworld_impress_title_selective_formatting_014
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_selective_formatting_014'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide data: (title, body_lines)
    slides_data = [
        (
            "Annual Report 2024",
            ["Meridian Technologies Inc.", "Fiscal Year Overview", "Prepared by: Office of the CFO"]
        ),
        (
            "Executive Summary",
            [
                "Revenue grew 18% year-over-year to $342M",
                "Operating income reached $67.4M, up from $51.2M",
                "Expanded into 4 new international markets",
                "Employee headcount increased from 1,840 to 2,210",
                "Launched 3 flagship product lines in Q3",
            ]
        ),
        (
            "Financial Highlights",
            [
                "Total Revenue: $342.1M (+18.3%)",
                "Gross Profit Margin: 54.7%",
                "EBITDA: $89.6M",
                "Net Income: $48.3M",
                "Earnings Per Share: $2.17",
                "Free Cash Flow: $55.9M",
            ]
        ),
        (
            "Regional Performance",
            [
                "North America: $178.4M — core market, stable growth",
                "Europe & MEA: $94.7M — accelerated adoption of cloud suite",
                "Asia-Pacific: $52.3M — entered Japan and South Korea",
                "Latin America: $16.7M — nascent market, pilot programs launched",
            ]
        ),
        (
            "Product & Innovation",
            [
                "MeridianCloud 4.0 — AI-assisted analytics dashboard",
                "SecureEdge v2 — zero-trust security framework",
                "DataBridge API — real-time enterprise integration layer",
                "R&D Investment: $41.2M (12% of revenue)",
                "Patents filed in 2024: 37",
            ]
        ),
        (
            "Operational Excellence",
            [
                "Customer satisfaction score: 4.6 / 5.0",
                "Support ticket resolution time reduced by 31%",
                "ISO 27001 recertification achieved in Q2",
                "Supply chain efficiency improved by 14%",
                "Operational cost savings: $8.7M",
            ]
        ),
        (
            "Outlook & Strategy 2025",
            [
                "Revenue target: $410M (projected +20%)",
                "Launch MeridianCloud 5.0 with generative AI features",
                "Acquire strategic partners in Southeast Asia",
                "Increase R&D spend to 15% of revenue",
                "Expand workforce to 2,600 employees globally",
            ]
        ),
    ]

    for idx, (title_text, body_lines) in enumerate(slides_data):
        if idx == 0:
            layout = prs.slide_layouts[0]  # Title Slide
        else:
            layout = prs.slide_layouts[1]  # Title + Content

        slide = prs.slides.add_slide(layout)

        # Set title — plain: black, regular, no underline
        title_shape = slide.shapes.title
        title_shape.text = title_text
        for para in title_shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = False
                run.font.underline = False
                run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
                run.font.size = Pt(36)

        # Set body content
        if idx == 0:
            # Title slide: use subtitle placeholder
            try:
                subtitle = slide.placeholders[1]
                subtitle.text = "\n".join(body_lines)
                for para in subtitle.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(20)
                        run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
            except (KeyError, IndexError):
                pass
        else:
            try:
                content = slide.placeholders[1]
                tf = content.text_frame
                tf.text = body_lines[0]
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(18)
                for line in body_lines[1:]:
                    p = tf.add_paragraph()
                    p.text = line
                    p.level = 0
                    for run in p.runs:
                        run.font.size = Pt(18)
            except (KeyError, IndexError):
                pass

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
