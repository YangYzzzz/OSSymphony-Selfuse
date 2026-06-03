"""
Reward Script: Build 8-slide Hamlet presentation with dark red/cream color scheme
Task ID: impress_stu_052
Domain: libreoffice_impress
Scoring:
  Component 1: Slide count == 8 (0.15)
  Component 2: Slide titles match (0.20)
  Component 3: Slide 2 has 5 bullet points (0.10)
  Component 4: Slide 3 has table with correct headers (0.15)
  Component 5: Slide 5 quotes are italic (0.10)
  Component 6: Slide 8 has 3 references (0.05)
  Component 7: Dark red background on all slides (0.15)
  Component 8: Cream text color used throughout (0.10)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_052'

EXPECTED_TITLES = [
    'Hamlet: Themes of Revenge and Mortality',
    'Plot Summary',
    'Major Characters',
    'Theme: Revenge',
    'Theme: Mortality',
    'Literary Devices',
    'Modern Relevance',
    'References',
]


def get_slide_title_text(slide):
    """Get the title text of a slide (first text shape, typically the heading)."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.paragraphs[0].text.strip()
            if text:
                return text
    return ""


def get_all_text_shapes(slide):
    """Get all text shapes recursively (handles groups)."""
    def extract(shape):
        results = []
        if hasattr(shape, "text") and hasattr(shape, "text_frame"):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.extend(extract(sub))
        return results
    out = []
    for shape in slide.shapes:
        out.extend(extract(shape))
    return out


def get_slide_bg_color(slide):
    """Get background color as hex string, handling inherited backgrounds."""
    fill = slide.background.fill
    if fill.type == 1:  # SOLID
        return str(fill.fore_color.rgb)
    elif fill.type == 5:  # inherited from master
        try:
            master_fill = slide.slide_layout.slide_master.background.fill
            if master_fill.type == 1:
                return str(master_fill.fore_color.rgb)
        except Exception:
            pass
    return None


def _check_slide_has_cream(text_shapes):
    """Return whether any text run on the slide uses cream (#FFFDD0) color."""
    for shape in text_shapes:
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                if (run.text or "").strip():
                    try:
                        if run.font.color.type is not None and str(run.font.color.rgb).upper() == 'FFFDD0':
                            return True  # derived from actual API check
                    except Exception:
                        pass
    return False


def verify_task(file_path):
    """Verify task completion with progressive scoring. Returns float 0.0-1.0."""
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Slide count == 8 (0.15 points)
    try:
        if num_slides == 8:
            print(f"PASS: Component 1 — Slide count is 8 (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — Expected 8 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide titles match expected titles (0.20 points)
    # Award proportional credit: 0.20 * (matched / 8)
    try:
        matched_titles = 0
        for idx, expected_title in enumerate(EXPECTED_TITLES):
            if idx < num_slides:
                actual_title = get_slide_title_text(prs.slides[idx])
                if expected_title.lower() in actual_title.lower():
                    matched_titles += 1
                else:
                    print(f"  Slide {idx+1} title mismatch: expected '{expected_title}', got '{actual_title}'")
        if matched_titles == 8:
            print(f"PASS: Component 2 — All 8 slide titles match (0.20 pts)")
            total_score += 0.20
        elif matched_titles > 0:
            pts = round(0.20 * matched_titles / 8, 3)
            print(f"PARTIAL: Component 2 — {matched_titles}/8 titles match ({pts} pts)")
            total_score += pts
        else:
            print(f"FAIL: Component 2 — No titles match")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 2 has 5 bullet points (0.10 points)
    try:
        if num_slides >= 2:
            slide2 = prs.slides[1]
            # Find the content text shape (not the title)
            text_shapes = get_all_text_shapes(slide2)
            bullet_count = 0
            for shape in text_shapes:
                title_text = get_slide_title_text(slide2)
                # Skip the title shape
                first_para_text = shape.text_frame.paragraphs[0].text.strip() if shape.text_frame.paragraphs else ""
                if first_para_text == title_text:
                    continue
                # Count non-empty paragraphs as bullets
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        bullet_count += 1
            if bullet_count >= 5:
                print(f"PASS: Component 3 — Slide 2 has {bullet_count} bullet points (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 3 — Slide 2 has {bullet_count} bullet points, expected >= 5")
        else:
            print(f"FAIL: Component 3 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 3 has a table with Character/Role/Significance headers (0.15 points)
    try:
        if num_slides >= 3:
            slide3 = prs.slides[2]
            table_obj = None
            for shape in slide3.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_obj = shape.table
                    break
            headers_match = (
                table_obj is not None
                and len(table_obj.columns) >= 3
                and len(table_obj.rows) >= 2
                and 'character' in table_obj.cell(0, 0).text.strip().lower()
                and 'role' in table_obj.cell(0, 1).text.strip().lower()
                and 'significance' in table_obj.cell(0, 2).text.strip().lower()
            )
            if table_obj is not None and headers_match:
                print(f"PASS: Component 4 — Slide 3 has table with correct headers (0.15 pts)")
                total_score += 0.15
            elif table_obj is not None:
                print(f"PARTIAL: Component 4 — Table found but headers incorrect (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 4 — No table found on slide 3")
        else:
            print(f"FAIL: Component 4 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slide 5 quotes are italic (0.10 points)
    try:
        if num_slides >= 5:
            slide5 = prs.slides[4]
            text_shapes = get_all_text_shapes(slide5)
            italic_quotes = 0
            total_quotes = 0
            for shape in text_shapes:
                title_text = get_slide_title_text(slide5)
                first_para_text = shape.text_frame.paragraphs[0].text.strip() if shape.text_frame.paragraphs else ""
                if first_para_text == title_text:
                    continue
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        total_quotes += 1
                        # Check if runs are italic
                        runs = [r for r in p.runs if (r.text or "").strip()]
                        if runs and all(r.font.italic is True for r in runs):
                            italic_quotes += 1
            if total_quotes >= 3 and italic_quotes == total_quotes:
                print(f"PASS: Component 5 — All {total_quotes} quotes on slide 5 are italic (0.10 pts)")
                total_score += 0.10
            elif italic_quotes > 0:
                pts = round(0.10 * italic_quotes / max(total_quotes, 1), 3)
                print(f"PARTIAL: Component 5 — {italic_quotes}/{total_quotes} quotes italic ({pts} pts)")
                total_score += pts
            else:
                print(f"FAIL: Component 5 — No italic quotes on slide 5 (found {total_quotes} quotes)")
        else:
            print(f"FAIL: Component 5 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Slide 8 has at least 3 references (0.05 points)
    try:
        if num_slides >= 8:
            slide8 = prs.slides[7]
            text_shapes = get_all_text_shapes(slide8)
            ref_count = 0
            for shape in text_shapes:
                title_text = get_slide_title_text(slide8)
                first_para_text = shape.text_frame.paragraphs[0].text.strip() if shape.text_frame.paragraphs else ""
                if first_para_text == title_text:
                    continue
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        ref_count += 1
            if ref_count >= 3:
                print(f"PASS: Component 6 — Slide 8 has {ref_count} references (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 6 — Slide 8 has {ref_count} references, expected >= 3")
        else:
            print(f"FAIL: Component 6 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Dark red (#8B0000) background on all slides (0.15 points)
    try:
        dark_red_count = 0
        for idx in range(num_slides):
            bg = get_slide_bg_color(prs.slides[idx])
            if bg is not None and bg.upper() == '8B0000':
                dark_red_count += 1
            else:
                print(f"  Slide {idx+1} bg color: {bg}")
        if num_slides >= 8 and dark_red_count == num_slides:
            print(f"PASS: Component 7 — All {num_slides} slides have dark red background (0.15 pts)")
            total_score += 0.15
        elif dark_red_count > 0:
            pts = round(0.15 * dark_red_count / max(num_slides, 1), 3)
            print(f"PARTIAL: Component 7 — {dark_red_count}/{num_slides} slides have dark red bg ({pts} pts)")
            total_score += pts
        else:
            print(f"FAIL: Component 7 — No slides have dark red background")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    # Component 8: Cream text color (#FFFDD0) used throughout (0.10 points)
    # Check that title runs across all slides use cream color
    try:
        cream_slides = 0
        for idx in range(num_slides):
            slide = prs.slides[idx]
            text_shapes = get_all_text_shapes(slide)
            cream_found_on_slide = _check_slide_has_cream(text_shapes)
            if cream_found_on_slide:
                cream_slides += 1
        if num_slides >= 8 and cream_slides == num_slides:
            print(f"PASS: Component 8 — All {num_slides} slides use cream text color (0.10 pts)")
            total_score += 0.10
        elif cream_slides > 0:
            pts = round(0.10 * cream_slides / max(num_slides, 1), 3)
            print(f"PARTIAL: Component 8 — {cream_slides}/{num_slides} slides use cream text ({pts} pts)")
            total_score += pts
        else:
            print(f"FAIL: Component 8 — No slides use cream text color")
    except Exception as e:
        print(f"ERROR: Component 8 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
persist_app_state()

if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
