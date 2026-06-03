"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 11’s heading is still in sentence case, and the design guidelines say every slide title has to be FULL UPPERCASE. How can I quickly flip that title on slide 11 to all caps in LibreOffice Impress?
Generated: 2025-09-10 12:07:56
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
import re
from pptx import Presentation

FILE_PATH = "/home/user/slide_11s_heading_is_still_in_sentence_case_and_the_design_guidelines_say_every_slide_title_has_to_b_golden.pptx"

def extract_slide_title(slide):
    """Return the most likely title text from a slide.

    Priority:
    1. Placeholder shapes marked as TITLE (type 1) or CENTER_TITLE (type 15)
    2. First text-bearing shape with non-empty text (fallback)
    """
    title_placeholders = []
    other_texts = []

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text.strip()
        if not text:
            continue

        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            if ph_type in (1, 15):  # 1 = TITLE, 15 = CENTER_TITLE
                title_placeholders.append(text)
        # Collect any text for fallback
        other_texts.append(text)

    if title_placeholders:
        return title_placeholders[0]
    return other_texts[0] if other_texts else ""

def is_all_caps(text):
    """Return True if *no* lowercase a-z characters appear in the text."""
    return re.search(r"[a-z]", text) is None

def verify_uppercase_title_on_slide_11(file_path):
    """Verify that slide 11 exists and its title is fully uppercase.

    Scoring (progressive):
    • 0.4  – Title text exists on slide 11
    • 0.6  – Title contains no lowercase letters (FULL UPPERCASE)
    Total possible: 1.0
    """
    max_score = 1.0
    score = 0.0

    print(f"Checking presentation: {file_path}")

    # --------------- Basic file checks (no points) ---------------
    if not os.path.exists(file_path):
        print("✗ File not found.")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        return 0.0

    slide_count = len(prs.slides)
    print(f"Presentation contains {slide_count} slides.")

    if slide_count < 11:
        print("✗ Presentation has fewer than 11 slides; cannot verify slide 11.")
        return 0.0

    # --------------- Retrieve and evaluate slide 11 ---------------
    slide_11 = prs.slides[10]  # zero-based index
    title_text = extract_slide_title(slide_11)

    if title_text:
        print(f"Found title on slide 11: \"{title_text}\"")
        score += 0.4  # Title exists
    else:
        print("✗ No title text found on slide 11.")
        return score  # Cannot award further points

    # --------------- Verify uppercase requirement ---------------
    if is_all_caps(title_text):
        print("✓ Title is fully uppercase.")
        score += 0.6
    else:
        print("✗ Title is not fully uppercase.")

    # --------------- Final score ---------------
    final_score = min(score, max_score)
    print(f"Total score: {final_score}/{max_score}")
    return final_score

if __name__ == "__main__":
    reward = verify_uppercase_title_on_slide_11(FILE_PATH)
    print(f"REWARD: {reward}")
