"""
FINAL REWARD SCRIPT - SUCCESS
Task: In LibreOffice Impress, how do I create a copy of Slide 1, move that duplicate to position 206, and then rename its title text to "Backup Title"?
Generated: 2025-09-10 19:23:46
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
import traceback
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

"""
Reward Script for LibreOffice Impress Task
Task:  
1. Duplicate Slide 1  
2. Move the duplicate so it becomes Slide 206 (1-indexed)  
3. Change the duplicate’s title to **"Backup Title"**

The script awards up to 1.0 points based on three ACTUAL verification steps:
• 0.2 – Slide 206 exists (shows the slide was created/moved)  
• 0.4 – Slide 206 title text exactly equals "Backup Title" (case–insensitive)  
• 0.4 – Content duplication check:
      – If Slide 1 has non-title text, the same text must appear on Slide 206 (full match ⇒ 0.4, partial ⇒ proportional)  
      – If Slide 1 has no non-title text, we fall back to comparing total shape counts; matching counts ⇒ 0.4

Returns a progressive score and prints detailed diagnostics plus the final
"REWARD: X.X" line as required.
"""

FILE_PATH = "/home/user/in_libreoffice_impress_how_do_i_create_a_copy_of_slide_1_move_that_duplicate_to_position_206_and_the_golden.pptx"

# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------

def _get_title_shape(slide):
    """Return the placeholder shape that acts as slide title if it exists."""
    for shp in slide.shapes:
        try:
            if shp.is_placeholder and shp.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                return shp
        except Exception:
            # Some shapes do not expose placeholder_format
            continue
    # Fallback: any shape whose name contains the word "title"
    for shp in slide.shapes:
        if hasattr(shp, "text") and shp.name and "title" in shp.name.lower():
            return shp
    return None


def _get_title_text(slide):
    title_shape = _get_title_shape(slide)
    if title_shape is not None and hasattr(title_shape, "text"):
        return title_shape.text.strip()
    return ""


def _get_non_title_texts(slide):
    """Collect non-empty texts from all shapes except the title placeholder."""
    title_shape = _get_title_shape(slide)
    texts = []
    for shp in slide.shapes:
        if shp == title_shape:
            continue
        if hasattr(shp, "text"):
            txt = shp.text.strip()
            if txt:
                texts.append(txt)
    return texts

# ---------------------------------------------------------------------------
# Main verification logic
# ---------------------------------------------------------------------------

def verify_impress_duplicate_task(pptx_path: str) -> float:
    max_score = 1.0
    score = 0.0

    try:
        # ---------- Basic file check (no points for simply existing) ----------
        if not os.path.exists(pptx_path):
            print(f"✗ File not found: {pptx_path}")
            return 0.0

        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        print(f"Loaded presentation with {slide_count} slides")

        # ---------- 1. Slide 206 existence (0.2) ----------
        target_idx = 205  # zero-based index for Slide 206
        if slide_count > target_idx:
            print("✓ Slide 206 exists (0.2 points)")
            score += 0.2
        else:
            print("✗ Slide 206 missing – cannot continue")
            return score  # Early exit; other checks depend on this slide

        slide1 = prs.slides[0]
        slide206 = prs.slides[target_idx]

        # ---------- 2. Title text verification (0.4) ----------
        title_text = _get_title_text(slide206)
        if title_text.lower() == "backup title":
            print('✓ Slide 206 title is "Backup Title" (0.4 points)')
            score += 0.4
        else:
            print(f'✗ Incorrect title on Slide 206 – found "{title_text}"')

        # ---------- 3. Content duplication verification (up to 0.4) ----------
        non_title_1 = _get_non_title_texts(slide1)
        non_title_206 = _get_non_title_texts(slide206)

        if non_title_1:
            matches = sum(1 for t in non_title_1 if t in non_title_206)
            ratio = matches / len(non_title_1)
            content_score = 0.4 * ratio
            if ratio == 1:
                print("✓ All non-title text duplicated (0.4 points)")
            else:
                print(f"✗ Only {matches}/{len(non_title_1)} non-title texts duplicated (\u2260 full match) – partial {content_score:.2f} points")
            score += content_score
        else:
            # Slide 1 has no non-title text – compare shape counts as a proxy
            shapes1 = len(slide1.shapes)
            shapes206 = len(slide206.shapes)
            if shapes1 == shapes206:
                print("✓ Shape count matches (layout duplicated) – awarding 0.4 points")
                score += 0.4
            else:
                print("✗ Shape counts differ – no duplication points")

    except Exception as err:
        print("✗ Exception during verification:", err)
        traceback.print_exc()
        return 0.0

    final_score = round(min(score, max_score), 2)
    print(f"Final Score: {final_score}")
    return final_score

# ---------------------------------------------------------------------------
# Execute verification and print reward
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    reward_value = verify_impress_duplicate_task(FILE_PATH)
    print(f"REWARD: {reward_value}")
