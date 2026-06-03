"""
Reward Script: Move the title placeholder on slide 1 to the bottom half of the slide.
Task ID: osworld_impress_title_position_bottom_003
Domain: libreoffice_impress
Scoring:
  Component 1 (0.6): Title textbox (TextBox 1) top edge is in the bottom half (top > slide_height/2)
  Component 2 (0.4): Title textbox is in the bottom half AND text content is preserved intact
"""

import os

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_position_bottom_003'

# Expected title text on slide 1 (ground truth from task exploration)
EXPECTED_TITLE_TEXT = 'NovaBrand Creative Pitch'
# Name of the title textbox on slide 1
TITLE_SHAPE_NAME = 'TextBox 1'


def persist_app_state():
    """Send Ctrl+S to persist any unsaved LibreOffice Impress edits."""
    try:
        import time
        os.environ["DISPLAY"] = ":0"
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0

    The task is to move the title textbox (TextBox 1) on slide 1 from the
    top half to the bottom half. Only the 'top' coordinate changes between
    initial_env (top=548640, in top half) and golden_env (top=4754880,
    in bottom half). Slide height = 6858000, midpoint = 3429000.

    Both scoring components are anchored to the position change, so they
    both fail on initial_env (top is not in bottom half) and pass on
    golden_env (top is in bottom half).
    """
    total_score = 0.0

    # Load presentation — if this fails, we cannot score anything
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition gate: slide 1 must exist
    if len(prs.slides) < 1:
        print("CRITICAL: Presentation has no slides.")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]
    slide_height = prs.slide_height
    vertical_midpoint = slide_height // 2

    print(f"INFO: slide_height={slide_height} EMU, vertical_midpoint={vertical_midpoint} EMU")

    # Find the title textbox by name on slide 1
    title_shape = None
    for shape in slide.shapes:
        if shape.name == TITLE_SHAPE_NAME:
            title_shape = shape
            break

    if title_shape is None:
        print(f"CRITICAL: Shape '{TITLE_SHAPE_NAME}' not found on slide 1. Cannot score.")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Found '{TITLE_SHAPE_NAME}' at top={title_shape.top}, height={title_shape.height}")

    # Component 1: Title textbox top edge is in the bottom half (0.6 points)
    # This checks the core task requirement: the title must be moved below the vertical midpoint.
    # FAILS on initial_env (top=548640 < 3429000, still in top half)
    # PASSES on golden_env  (top=4754880 > 3429000, moved to bottom half)
    try:
        title_top = title_shape.top
        if title_top > vertical_midpoint:
            print(f"PASS: Component 1 — Title top ({title_top}) > midpoint ({vertical_midpoint}): title is in the bottom half (0.6 pts)")
            total_score += 0.6
        else:
            print(f"FAIL: Component 1 — Title top ({title_top}) is NOT above midpoint ({vertical_midpoint}): still in top half")
    except Exception as e:
        print(f"ERROR: Component 1 — Could not check title position: {e}")

    # Component 2: Title is in the bottom half AND text is preserved intact (0.4 points)
    # Compound check: both position (same as Component 1) AND text integrity.
    # Anchored to the position change — this FAILS on initial_env because the position
    # condition is NOT met there (title is still in top half).
    # FAILS on initial_env (top not in bottom half → condition fails at first check)
    # PASSES on golden_env  (top in bottom half AND text == 'NovaBrand Creative Pitch')
    try:
        title_top = title_shape.top
        actual_text = title_shape.text_frame.text.strip() if title_shape.has_text_frame else ""
        if title_top > vertical_midpoint and actual_text == EXPECTED_TITLE_TEXT:
            print(f"PASS: Component 2 — Title is in bottom half AND text is '{actual_text}' (0.4 pts)")
            total_score += 0.4
        elif title_top <= vertical_midpoint:
            print(f"FAIL: Component 2 — Title not yet in bottom half (top={title_top}), compound check skipped")
        else:
            print(f"FAIL: Component 2 — Title is in bottom half but text mismatch: expected '{EXPECTED_TITLE_TEXT}', found '{actual_text}'")
    except Exception as e:
        print(f"ERROR: Component 2 — Could not perform compound check: {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
persist_app_state()
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
