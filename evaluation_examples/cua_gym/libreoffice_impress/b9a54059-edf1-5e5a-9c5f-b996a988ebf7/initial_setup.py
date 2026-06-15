"""
Initial Setup: Move title placeholder on slide 1 to the bottom half of the slide.
Task ID: osworld_impress_title_position_bottom_003
Domain: libreoffice_impress

Creates a 4-slide creative pitch deck with slide 1 having:
- A dark-colored background (simulating a background image style)
- A title placeholder positioned in the TOP half of the slide
The agent's task is to move this title to the bottom half.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_position_bottom_003'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard widescreen 16:9 slide dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # ---------------------------------------------------------------
    # Slide 1: Cover slide — title in TOP half, dark creative bg
    # ---------------------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Dark gradient-style background
    bg_fill = slide1.background.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # Title textbox positioned in the TOP half (top ~10% of slide)
    # Vertical midpoint = slide_height / 2 = Inches(3.75)
    # Title top = Inches(0.6) — clearly in top half
    title_left = Inches(1.0)
    title_top = Inches(0.6)     # TOP half — must be < Inches(3.75)
    title_width = Inches(11.333)
    title_height = Inches(1.5)

    title_box = slide1.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = False
    title_p = title_tf.paragraphs[0]
    title_p.alignment = PP_ALIGN.LEFT
    title_run = title_p.add_run()
    title_run.text = "NovaBrand Creative Pitch"
    title_run.font.name = "Calibri"
    title_run.font.size = Pt(44)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Subtitle / tagline on slide 1
    sub_left = Inches(1.0)
    sub_top = Inches(2.3)
    sub_width = Inches(8.0)
    sub_height = Inches(0.8)
    sub_box = slide1.shapes.add_textbox(sub_left, sub_top, sub_width, sub_height)
    sub_tf = sub_box.text_frame
    sub_p = sub_tf.paragraphs[0]
    sub_p.alignment = PP_ALIGN.LEFT
    sub_run = sub_p.add_run()
    sub_run.text = "Redefining Visual Storytelling · 2025"
    sub_run.font.name = "Calibri"
    sub_run.font.size = Pt(20)
    sub_run.font.color.rgb = RGBColor(0xA0, 0xC4, 0xFF)

    # Decorative accent bar
    from pptx.util import Emu as E
    accent = slide1.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE rectangle
        Inches(1.0), Inches(2.1), Inches(6.0), Emu(45720)  # thin bar
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(0x6C, 0x63, 0xFF)
    accent.line.fill.background()

    # ---------------------------------------------------------------
    # Slide 2: About Us
    # ---------------------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    bg2 = slide2.background.fill
    bg2.solid()
    bg2.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

    h2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1.0))
    h2.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    r2h = h2.text_frame.paragraphs[0].add_run()
    r2h.text = "About NovaBrand"
    r2h.font.name = "Calibri"
    r2h.font.size = Pt(36)
    r2h.font.bold = True
    r2h.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    body2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.5))
    body2.text_frame.word_wrap = True
    bullets2 = [
        "Founded in 2018 by a team of award-winning designers and strategists",
        "Offices in San Francisco, Berlin, and Singapore",
        "Trusted by 120+ global brands across technology, fashion, and lifestyle sectors",
        "Recognized by Fast Company as one of the 'Most Creative Companies' in 2024",
        "Our philosophy: bold ideas, rigorous craft, measurable impact",
    ]
    for i, b in enumerate(bullets2):
        if i == 0:
            p = body2.text_frame.paragraphs[0]
        else:
            p = body2.text_frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"• {b}"
        run.font.name = "Calibri"
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ---------------------------------------------------------------
    # Slide 3: Our Services
    # ---------------------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    bg3 = slide3.background.fill
    bg3.solid()
    bg3.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    h3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1.0))
    h3.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    r3h = h3.text_frame.paragraphs[0].add_run()
    r3h.text = "What We Deliver"
    r3h.font.name = "Calibri"
    r3h.font.size = Pt(36)
    r3h.font.bold = True
    r3h.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    services = [
        ("Brand Identity", "Logos, colour systems, typography, brand guidelines"),
        ("Campaign Strategy", "Multi-channel narratives and audience engagement planning"),
        ("Motion & Film", "TV commercials, social reels, branded documentary series"),
        ("Digital Experience", "Website design, interactive microsites, app UI/UX"),
    ]
    for idx, (svc, desc) in enumerate(services):
        col = idx % 2
        row = idx // 2
        bx = slide3.shapes.add_textbox(
            Inches(0.8 + col * 6.2), Inches(1.8 + row * 2.4),
            Inches(5.6), Inches(2.0)
        )
        bx.text_frame.word_wrap = True
        p_svc = bx.text_frame.paragraphs[0]
        r_svc = p_svc.add_run()
        r_svc.text = svc
        r_svc.font.name = "Calibri"
        r_svc.font.size = Pt(20)
        r_svc.font.bold = True
        r_svc.font.color.rgb = RGBColor(0x6C, 0x63, 0xFF)
        p_desc = bx.text_frame.add_paragraph()
        r_desc = p_desc.add_run()
        r_desc.text = desc
        r_desc.font.name = "Calibri"
        r_desc.font.size = Pt(15)
        r_desc.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    # ---------------------------------------------------------------
    # Slide 4: Call to Action / Contact
    # ---------------------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    bg4 = slide4.background.fill
    bg4.solid()
    bg4.fore_color.rgb = RGBColor(0x6C, 0x63, 0xFF)

    h4 = slide4.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.0), Inches(1.4))
    h4.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r4h = h4.text_frame.paragraphs[0].add_run()
    r4h.text = "Let's Create Something Extraordinary"
    r4h.font.name = "Calibri"
    r4h.font.size = Pt(38)
    r4h.font.bold = True
    r4h.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    contact_box = slide4.shapes.add_textbox(Inches(3.5), Inches(3.8), Inches(6.0), Inches(1.5))
    contact_box.text_frame.word_wrap = True
    contact_lines = [
        "hello@novabrand.studio",
        "+1 (415) 823-9910  |  novabrand.studio",
    ]
    for i, line in enumerate(contact_lines):
        if i == 0:
            cp = contact_box.text_frame.paragraphs[0]
        else:
            cp = contact_box.text_frame.add_paragraph()
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = line
        cr.font.name = "Calibri"
        cr.font.size = Pt(18)
        cr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'  Slide 1 title top position: {title_top} EMU  (slide_height/2 = {int(slide_height/2)} EMU)')
    print(f'  Title is in the TOP half: {title_top < slide_height / 2}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
