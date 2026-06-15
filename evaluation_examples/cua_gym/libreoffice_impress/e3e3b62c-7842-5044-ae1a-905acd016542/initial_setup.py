"""
Initial Setup: Create Senior Project presentation with 8 slides, slide 4 empty (Budget Overview)
Task ID: impress_stu_093
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_093'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    if body_lines:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, line in enumerate(body_lines):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                p = tf.add_paragraph()
                p.text = line
    return slide


def add_title_only_slide(prs, title_text):
    # Use layout 5 (blank) and add title manually for a clean empty slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(prs, "Senior Project: Smart Campus Navigation",
                    "Team Horizon  |  Spring 2026  |  Dr. Emily Foster, Advisor")

    # Slide 2: Project Overview
    add_content_slide(prs, "Project Overview", [
        "Develop an indoor navigation system for the university campus",
        "Utilize Bluetooth Low Energy beacons for precise positioning",
        "Mobile app with real-time wayfinding and accessibility features",
        "Integration with campus event calendar and room booking system",
        "Target accuracy: sub-2-meter positioning in all buildings",
    ])

    # Slide 3: Team Members
    add_content_slide(prs, "Team Members", [
        "Aisha Patel - Project Lead & Backend Developer",
        "Ryan Nakamura - Mobile App Developer",
        "Sophia Rodriguez - UX/UI Designer",
        "Liam O'Brien - Hardware & Beacon Deployment",
        "Mei-Lin Chang - Data Analytics & Testing",
    ])

    # Slide 4: Budget Overview (EMPTY - no chart, no data)
    add_title_only_slide(prs, "Budget Overview")

    # Slide 5: Timeline
    add_content_slide(prs, "Project Timeline", [
        "Phase 1 (Jan-Feb): Requirements gathering and beacon procurement",
        "Phase 2 (Mar-Apr): Beacon installation and initial calibration",
        "Phase 3 (Apr-May): Mobile app development and integration",
        "Phase 4 (May-Jun): User testing and iteration",
        "Final Presentation: June 15, 2026",
    ])

    # Slide 6: Methodology
    add_content_slide(prs, "Methodology", [
        "Agile development with 2-week sprint cycles",
        "Fingerprinting-based localization using BLE signal strength",
        "A/B testing for navigation UI with 50+ student participants",
        "Continuous integration pipeline with automated regression tests",
        "Weekly stakeholder demos for iterative feedback",
    ])

    # Slide 7: Results
    add_content_slide(prs, "Preliminary Results", [
        "Beacon deployment completed in 3 campus buildings",
        "Average positioning accuracy: 1.7 meters (exceeds target)",
        "Mobile app prototype tested by 25 students in pilot group",
        "User satisfaction rating: 4.3 / 5.0 in initial survey",
        "Reduced average time to find classrooms by 40%",
    ])

    # Slide 8: Conclusion
    add_content_slide(prs, "Conclusion & Next Steps", [
        "Core navigation system proven effective in pilot deployment",
        "Expand beacon network to remaining 5 campus buildings",
        "Add accessibility features: screen reader support, wheelchair routes",
        "Explore integration with campus security and emergency systems",
        "Open-source release planned for Fall 2026",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
