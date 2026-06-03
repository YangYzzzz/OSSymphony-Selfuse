"""
Reward Script: Doughnut chart on slide 4 with budget segments, colors, data labels, and center label
Task ID: impress_stu_093
Domain: libreoffice_impress
Scoring:
  Component 1: Doughnut chart exists on slide 4 (0.20)
  Component 2: Chart title is 'Budget Breakdown' (0.10)
  Component 3: Correct categories and values (0.25)
  Component 4: Correct segment colors (0.20)
  Component 5: Data labels show category names and percentages (0.15)
  Component 6: Center total label '$1,000' text box on slide 4 (0.10)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_093'

EXPECTED_CATEGORIES = ['Materials', 'Software Licenses', 'Printing', 'Travel', 'Miscellaneous']
EXPECTED_VALUES = [450.0, 200.0, 150.0, 100.0, 100.0]
EXPECTED_COLORS = ['3498DB', '2ECC71', 'E67E22', '9B59B6', '95A5A6']


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # slide 4 (0-indexed)

    # Find chart shape on slide 4
    chart_shape = None
    for shape in slide.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Doughnut chart exists on slide 4 (0.20 points)
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            # Check it's a doughnut chart (enum value -4120)
            if chart.chart_type == -4120:
                print(f"PASS: Component 1 — Doughnut chart found on slide 4 (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 1 — Chart found but type is {chart.chart_type}, expected DOUGHNUT (-4120)")
        else:
            print(f"FAIL: Component 1 — No chart found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    if chart_shape is None:
        # No chart means no further checks possible
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    chart = chart_shape.chart

    # Component 2: Chart title is 'Budget Breakdown' (0.10 points)
    try:
        if chart.has_title:
            title_text = chart.chart_title.text_frame.text.strip()
            if title_text == 'Budget Breakdown':
                print(f"PASS: Component 2 — Chart title is 'Budget Breakdown' (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 2 — Chart title is '{title_text}', expected 'Budget Breakdown'")
        else:
            print(f"FAIL: Component 2 — Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Correct categories and values (0.25 points)
    try:
        plot = chart.plots[0]
        categories = [str(c) for c in plot.categories]
        series_values = tuple(chart.series[0].values)

        cats_match = categories == EXPECTED_CATEGORIES
        vals_match = len(series_values) == len(EXPECTED_VALUES) and all(
            abs(a - b) < 0.01 for a, b in zip(series_values, EXPECTED_VALUES)
        )

        if cats_match and vals_match:
            print(f"PASS: Component 3 — Categories and values match (0.25 pts)")
            total_score += 0.25
        else:
            if not cats_match:
                print(f"FAIL: Component 3 — Categories: {categories}, expected {EXPECTED_CATEGORIES}")
            if not vals_match:
                print(f"FAIL: Component 3 — Values: {series_values}, expected {EXPECTED_VALUES}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Correct segment colors (0.20 points)
    try:
        series = chart.series[0]
        color_matches = 0
        for idx, expected_color in enumerate(EXPECTED_COLORS):
            try:
                pt = series.points[idx]
                fill = pt.format.fill
                if fill.type is not None:
                    actual_color = str(fill.fore_color.rgb).upper()
                    if actual_color == expected_color.upper():
                        color_matches += 1
                    else:
                        print(f"  Color mismatch point {idx}: {actual_color} vs expected {expected_color}")
                else:
                    print(f"  Point {idx} has no solid fill")
            except Exception as e:
                print(f"  Point {idx} color check error: {e}")

        if color_matches == 5:
            print(f"PASS: Component 4 — All 5 segment colors correct (0.20 pts)")
            total_score += 0.20
        elif color_matches >= 3:
            partial = round(0.20 * color_matches / 5, 2)
            print(f"PARTIAL: Component 4 — {color_matches}/5 colors correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — Only {color_matches}/5 colors correct")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Data labels show category names and percentages (0.15 points)
    try:
        plot = chart.plots[0]
        if plot.has_data_labels:
            dl = plot.data_labels
            show_cat = dl.show_category_name
            show_pct = dl.show_percentage

            if show_cat and show_pct:
                print(f"PASS: Component 5 — Data labels show category names and percentages (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 5 — show_category_name={show_cat}, show_percentage={show_pct}")
        else:
            print(f"FAIL: Component 5 — No data labels on chart")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Center total label '$1,000' on slide 4 (0.10 points)
    try:
        label_texts = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame
            and not (hasattr(shape, 'has_chart') and shape.has_chart)
        ]
        matching = [t for t in label_texts if '$1,000' in t or '$1000' in t]

        if len(matching) > 0:
            print(f"PASS: Component 6 — Center label '$1,000' found on slide 4 (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 6 — No '$1,000' text box found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
