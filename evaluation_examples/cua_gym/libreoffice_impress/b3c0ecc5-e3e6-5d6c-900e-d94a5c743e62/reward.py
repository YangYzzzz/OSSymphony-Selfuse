"""
Reward Script: Set font to Arial and font size to 20pt for all text on slides 1, 2, and 3
Task ID: osworld_impress_global_font_change_007
Domain: libreoffice_impress
Scoring:
  Component 1: All text runs on slide 1 use Arial font (0.30 pts)
  Component 2: All text runs on slides 2 and 3 use Arial font (0.40 pts)
  Component 3: All text runs on slides 1, 2, and 3 use 20pt font size (0.30 pts)
  Total: 1.0

Notes:
  - 20pt = 254000 EMU (20 * 12700)
  - Target font: Arial
  - Scope: slides 1, 2, 3 only (0-indexed: 0, 1, 2)
  - Slide 4 (0-indexed 3) was already Arial but at different sizes — not in scope
  - Slide 5 (0-indexed 4) uses Calibri — not in scope
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_global_font_change_007'

# 20pt in EMU = 20 * 12700 = 254000
TARGET_FONT_NAME = 'Arial'
TARGET_FONT_SIZE_EMU = 20 * 12700  # 254000


def get_all_runs_on_slide(slide):
    """Return a list of (shape_name, para_idx, run) for all non-empty text runs on a slide."""
    results = []

    def collect(shape):
        if shape.has_text_frame:
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                for run in para.runs:
                    if (run.text or '').strip():
                        results.append((shape.name, para_idx, run))
        # Recurse into group shapes
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                collect(sub)

    for shape in slide.shapes:
        collect(shape)
    return results


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"CRITICAL: Expected at least 3 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # -----------------------------------------------------------------------
    # Component 1: All text runs on slide 1 use Arial font (0.30 pts)
    # Fails on initial (Calibri+Georgia), passes on golden (all Arial)
    # -----------------------------------------------------------------------
    try:
        slide1_runs = get_all_runs_on_slide(prs.slides[0])
        if len(slide1_runs) == 0:
            print("FAIL: Component 1 — slide 1 has no non-empty text runs")
        else:
            non_arial_slide1 = [(sn, pi, r.text[:40], r.font.name)
                                for sn, pi, r in slide1_runs
                                if r.font.name != TARGET_FONT_NAME]
            if not non_arial_slide1:
                print(f"PASS: Component 1 — all {len(slide1_runs)} run(s) on slide 1 use Arial (0.30 pts)")
                total_score += 0.30
            else:
                print(f"FAIL: Component 1 — {len(non_arial_slide1)} run(s) on slide 1 are NOT Arial:")
                for sn, pi, txt, fn in non_arial_slide1:
                    print(f"       shape={sn}, para={pi}, text={repr(txt)}, font={fn}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # -----------------------------------------------------------------------
    # Component 2: All text runs on slides 2 and 3 use Arial font (0.40 pts)
    # Fails on initial (Georgia+Calibri mix), passes on golden (all Arial)
    # -----------------------------------------------------------------------
    try:
        slide2_runs = get_all_runs_on_slide(prs.slides[1])
        slide3_runs = get_all_runs_on_slide(prs.slides[2])

        non_arial_2_3 = []
        for slide_idx, slide_runs in [(2, slide2_runs), (3, slide3_runs)]:
            for sn, pi, r in slide_runs:
                if r.font.name != TARGET_FONT_NAME:
                    non_arial_2_3.append((slide_idx, sn, pi, r.text[:40], r.font.name))

        total_runs_2_3 = len(slide2_runs) + len(slide3_runs)
        if total_runs_2_3 == 0:
            print("FAIL: Component 2 — slides 2 and 3 have no non-empty text runs")
        elif not non_arial_2_3:
            print(f"PASS: Component 2 — all {total_runs_2_3} run(s) on slides 2 and 3 use Arial (0.40 pts)")
            total_score += 0.40
        else:
            print(f"FAIL: Component 2 — {len(non_arial_2_3)} run(s) on slides 2/3 are NOT Arial:")
            for slide_idx, sn, pi, txt, fn in non_arial_2_3:
                print(f"       slide={slide_idx}, shape={sn}, para={pi}, text={repr(txt)}, font={fn}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # -----------------------------------------------------------------------
    # Component 3: All text runs on slides 1, 2, and 3 use 20pt font size (0.30 pts)
    # Fails on initial (various sizes: 36pt, 28pt, 18pt, 16pt, 17pt, etc.)
    # Passes on golden (all 20pt = 254000 EMU)
    # -----------------------------------------------------------------------
    try:
        all_runs_1_3 = (get_all_runs_on_slide(prs.slides[0]) +
                        get_all_runs_on_slide(prs.slides[1]) +
                        get_all_runs_on_slide(prs.slides[2]))

        if len(all_runs_1_3) == 0:
            print("FAIL: Component 3 — slides 1-3 have no non-empty text runs")
        else:
            wrong_size = []
            for sn, pi, r in all_runs_1_3:
                if r.font.size != TARGET_FONT_SIZE_EMU:
                    size_pt = round(r.font.size / 12700, 1) if r.font.size else None
                    wrong_size.append((sn, pi, r.text[:40], r.font.size, size_pt))

            if not wrong_size:
                print(f"PASS: Component 3 — all {len(all_runs_1_3)} run(s) on slides 1-3 use 20pt ({TARGET_FONT_SIZE_EMU} EMU) (0.30 pts)")
                total_score += 0.30
            else:
                print(f"FAIL: Component 3 — {len(wrong_size)} run(s) on slides 1-3 are NOT 20pt:")
                for sn, pi, txt, sz, sz_pt in wrong_size:
                    print(f"       shape={sn}, para={pi}, text={repr(txt)}, size={sz} ({sz_pt}pt)")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 4)
    print(f"\nScore: {round(total_score, 4)}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in a given env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
