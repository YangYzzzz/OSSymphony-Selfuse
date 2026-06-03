"""
Initial Setup: Sales pitch presentation with mixed fonts on slides 1-3
Task ID: osworld_impress_global_font_change_007
Domain: libreoffice_impress

Creates a 5-slide sales pitch deck with Calibri and Georgia fonts at various
sizes on slides 1-3. Slides 4-5 use consistent formatting.
The agent task is to change all text on slides 1-3 to Arial at 20pt.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'  # VM path — all scripts run on the VM
TASK_ID = 'osworld_impress_global_font_change_007'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_run_font(run, font_name, font_size_pt, bold=False, color_rgb=None):
    """Helper to set font name, size, bold on a run."""
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    if color_rgb:
        run.font.color.rgb = color_rgb


def add_text_to_shape(shape, text, font_name, font_size_pt, bold=False,
                       alignment=PP_ALIGN.LEFT, color_rgb=None):
    """Add text to a shape's text frame, clearing existing content."""
    tf = shape.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = alignment
    run = para.add_run()
    run.text = text
    set_run_font(run, font_name, font_size_pt, bold=bold, color_rgb=color_rgb)


def create_initial():
    prs = Presentation()
    # Standard 16:9 widescreen
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # -------------------------------------------------------------------------
    # Slide 1: Title Slide — mixed fonts (Calibri title, Georgia subtitle)
    # -------------------------------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout

    # Title placeholder
    title1 = slide1.shapes.title
    title1.text = ''
    title1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run_t1 = title1.text_frame.paragraphs[0].add_run()
    run_t1.text = 'NextGen Solutions: Driving Your Business Forward'
    set_run_font(run_t1, 'Calibri', 36, bold=True,
                 color_rgb=RGBColor(0x1F, 0x39, 0x64))

    # Subtitle placeholder
    sub1 = slide1.placeholders[1]
    sub1.text = ''
    sub1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run_s1 = sub1.text_frame.paragraphs[0].add_run()
    run_s1.text = 'A Comprehensive Sales Proposal — Q2 2025'
    set_run_font(run_s1, 'Georgia', 20, bold=False,
                 color_rgb=RGBColor(0x40, 0x40, 0x40))

    # -------------------------------------------------------------------------
    # Slide 2: Company Overview — Georgia title, Calibri body at mixed sizes
    # -------------------------------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content

    title2 = slide2.shapes.title
    title2.text = ''
    title2.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    run_t2 = title2.text_frame.paragraphs[0].add_run()
    run_t2.text = 'About NextGen Solutions'
    set_run_font(run_t2, 'Georgia', 28, bold=True,
                 color_rgb=RGBColor(0x1F, 0x39, 0x64))

    content2 = slide2.placeholders[1]
    content2.text = ''
    tf2 = content2.text_frame
    tf2.word_wrap = True

    bullets2 = [
        ('Founded in 2010, headquartered in San Francisco, CA', 'Calibri', 18),
        ('Over 500 enterprise clients across 30 countries', 'Calibri', 16),
        ('Award-winning platform with 99.9% uptime guarantee', 'Georgia', 17),
        ('Dedicated team of 1,200+ engineers and consultants', 'Calibri', 18),
        ('Annual revenue exceeding $250M in fiscal year 2024', 'Georgia', 16),
    ]

    for i, (text, font, size) in enumerate(bullets2):
        if i == 0:
            para = tf2.paragraphs[0]
        else:
            para = tf2.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.level = 0
        run = para.add_run()
        run.text = text
        set_run_font(run, font, size)

    # -------------------------------------------------------------------------
    # Slide 3: Value Proposition — mixed Calibri/Georgia at various sizes
    # -------------------------------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content

    title3 = slide3.shapes.title
    title3.text = ''
    title3.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    run_t3 = title3.text_frame.paragraphs[0].add_run()
    run_t3.text = 'Why Choose Us?'
    set_run_font(run_t3, 'Calibri', 30, bold=True,
                 color_rgb=RGBColor(0x1F, 0x39, 0x64))

    content3 = slide3.placeholders[1]
    content3.text = ''
    tf3 = content3.text_frame
    tf3.word_wrap = True

    bullets3 = [
        ('Reduce operational costs by up to 35% within 12 months', 'Georgia', 17),
        ('Seamless integration with your existing ERP and CRM systems', 'Calibri', 18),
        ('AI-powered analytics delivering real-time business insights', 'Georgia', 16),
        ('24/7 premium support with dedicated account managers', 'Calibri', 17),
        ('Scalable pricing tailored to your growth trajectory', 'Georgia', 18),
        ('Proven ROI: average payback period of 8 months', 'Calibri', 16),
    ]

    for i, (text, font, size) in enumerate(bullets3):
        if i == 0:
            para = tf3.paragraphs[0]
        else:
            para = tf3.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.level = 0
        run = para.add_run()
        run.text = text
        set_run_font(run, font, size)

    # -------------------------------------------------------------------------
    # Slide 4: Case Studies — consistent Arial formatting (NOT affected by task)
    # -------------------------------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])

    title4 = slide4.shapes.title
    title4.text = ''
    title4.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    run_t4 = title4.text_frame.paragraphs[0].add_run()
    run_t4.text = 'Customer Success Stories'
    set_run_font(run_t4, 'Arial', 28, bold=True,
                 color_rgb=RGBColor(0x1F, 0x39, 0x64))

    content4 = slide4.placeholders[1]
    content4.text = ''
    tf4 = content4.text_frame
    tf4.word_wrap = True

    bullets4 = [
        'TechCorp Inc. — 42% reduction in IT overhead costs',
        'GlobalBank Ltd. — Deployed across 18 countries in 90 days',
        'HealthFirst Hospital — $3.2M saved annually in admin costs',
        'RetailMax Group — 60% improvement in supply chain efficiency',
    ]
    for i, text in enumerate(bullets4):
        if i == 0:
            para = tf4.paragraphs[0]
        else:
            para = tf4.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = text
        set_run_font(run, 'Arial', 18)

    # -------------------------------------------------------------------------
    # Slide 5: Next Steps / Call to Action — consistent Calibri formatting
    # -------------------------------------------------------------------------
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])

    title5 = slide5.shapes.title
    title5.text = ''
    title5.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run_t5 = title5.text_frame.paragraphs[0].add_run()
    run_t5.text = 'Your Next Steps with NextGen'
    set_run_font(run_t5, 'Calibri', 28, bold=True,
                 color_rgb=RGBColor(0x1F, 0x39, 0x64))

    content5 = slide5.placeholders[1]
    content5.text = ''
    tf5 = content5.text_frame
    tf5.word_wrap = True

    bullets5 = [
        'Schedule a 30-minute discovery call with our solutions team',
        'Request a live product demo customised to your industry',
        'Review the tailored commercial proposal (attached)',
        'Sign the agreement and launch your onboarding in under 2 weeks',
    ]
    for i, text in enumerate(bullets5):
        if i == 0:
            para = tf5.paragraphs[0]
        else:
            para = tf5.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = text
        set_run_font(run, 'Calibri', 18)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup — open the file in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
