"""
FINAL REWARD SCRIPT - SUCCESS
Task: Add page numbers to the footer on pages 2–5 only, centered.
Generated: 2025-10-17 13:12:12
Status: success
Model: azure-o3
Total Steps: 7
"""

import os
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

"""
Reward Script: Verify that page numbers were added ONLY to slides 2–5, they are
located in the footer area (any text containing the PowerPoint slide‐number
field "<#>") and that the text is centre-aligned.  Slides 1 and 6+ must NOT
contain a page-number field.  Progressive scoring is based on:
  • Correct presence / absence on every slide  (60 %)
  • Correct centre alignment on slides that must contain page numbers (40 %)
The script prints detailed diagnostics and finishes with  "REWARD: X.X"  where
X.X ∈ [0.0, 1.0].
"""

SLIDE_NUMBER_PH_TYPE = 6  # placeholder_format.type constant for slide numbers


def _page_number_shapes(slide):
    """Return a list of shapes that represent a page-number field."""
    shapes = []
    for shape in slide.shapes:
        # 1) Native slide-number placeholder
        if shape.is_placeholder:
            try:
                if shape.placeholder_format.type == SLIDE_NUMBER_PH_TYPE:
                    shapes.append(shape)
                    continue
            except Exception:
                pass  # malformed placeholder – ignore
        # 2) Text box that contains the field code "<#>"
        if getattr(shape, "has_text_frame", False):
            tf = shape.text_frame
            if tf and "<#>" in tf.text:
                shapes.append(shape)
    return shapes


def _is_center_aligned(shape):
    """Check if first paragraph of a shape is centre-aligned."""
    if getattr(shape, "has_text_frame", False) and shape.text_frame.paragraphs:
        return shape.text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER
    return False


def verify_page_numbers(file_path: str) -> float:
    if not os.path.exists(file_path):
        print("✗ File not found:", file_path)
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print("✗ Unable to open PPTX:", e)
        print("REWARD: 0.0")
        return 0.0

    total_slides = len(prs.slides)
    print(f"Loaded presentation with {total_slides} slides\n")

    # Scoring counters
    presence_correct = 0  # how many slides have correct presence/absence
    alignment_correct = 0  # slides 2-5 that are correctly centred
    alignment_possible = 0  # number of slides that should contain page number (2-5)

    for idx, slide in enumerate(prs.slides, start=1):
        shapes = _page_number_shapes(slide)
        found = bool(shapes)
        expected = 2 <= idx <= 5  # requirement states only slides 2–5

        # Presence / absence verification
        if expected == found:
            presence_correct += 1
        print(f"Slide {idx}: expected_page_number={expected}, found={found}")

        # Alignment verification for slides that must have page numbers
        if expected:
            alignment_possible += 1
            if found and any(_is_center_aligned(s) for s in shapes):
                alignment_correct += 1
                print("  ✓ Centre alignment verified")
            else:
                print("  ✗ Centre alignment incorrect or missing")
        print()

    # Progressive scoring
    presence_score = 0.6 * (presence_correct / total_slides) if total_slides else 0
    alignment_score = (
        0.4 * (alignment_correct / alignment_possible) if alignment_possible else 0
    )

    final_score = round(min(presence_score + alignment_score, 1.0), 3)

    print(f"Presence score:  {presence_score:.3f}  (max 0.6)")
    print(f"Alignment score: {alignment_score:.3f}  (max 0.4)")
    print(f"TOTAL SCORE: {final_score}")
    print(f"REWARD: {final_score}")

    return final_score


# ------------------------- RUN VERIFICATION ------------------------------
FILE_PATH = "/home/user/add_page_numbers_to_the_footer_on_pages_25_only_centered.pptx"
verify_page_numbers(FILE_PATH)
