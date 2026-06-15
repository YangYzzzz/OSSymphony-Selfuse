"""
Reward Script: Two-tone split background title slide
Task ID: impress_rp_001
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Left-half rectangle with dark teal (#004D4D) fill
  Component 2 (0.35): Title 'Innovation Summit 2026' — 40pt bold Montserrat white
  Component 3 (0.30): Subtitle 'Shaping Tomorrow Together' — 22pt italic Montserrat #004D4D
  Component 4 (0.10): Spatial layout — title on left half, subtitle on right half
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_001'


def persist_app_state(domain: str):
    """Save any unsaved GUI edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Pt, Inches, Emu
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 1 slide
    if len(prs.slides) == 0:
        print("FAIL: Presentation has no slides")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]
    slide_width = prs.slide_width  # ~12191695 EMU = 13.33 inches
    slide_height = prs.slide_height  # 6858000 EMU = 7.5 inches
    midpoint_x = slide_width / 2

    # Collect shapes by type for analysis
    rectangles = []
    textboxes = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or shape.shape_type == 1:
            rectangles.append(shape)
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                textboxes.append(shape)

    # =========================================================================
    # Component 1: Left-half rectangle with #004D4D fill (0.25 points)
    # =========================================================================
    try:
        teal_rect_found = False
        for rect in rectangles:
            try:
                fill = rect.fill
                if fill.type is not None and fill.type == 1:  # SOLID fill
                    color_rgb = str(fill.fore_color.rgb).upper()
                    if color_rgb == "004D4D":
                        # Check it covers approximately the left half
                        # Should start near x=0 and be roughly half the slide width
                        left_ok = rect.left <= Inches(0.5)  # starts near left edge
                        width_ok = abs(rect.width - midpoint_x) / midpoint_x <= 0.15  # ~half width
                        height_ok = abs(rect.height - slide_height) / slide_height <= 0.1  # ~full height
                        if left_ok and width_ok and height_ok:
                            teal_rect_found = True
                            print(f"PASS: Component 1 — Teal rectangle found: "
                                  f"left={rect.left}, width={rect.width}, height={rect.height}, "
                                  f"fill=#{color_rgb} (0.25 pts)")
                            total_score += 0.25
                            break
            except Exception:
                continue

        if not teal_rect_found:
            print(f"FAIL: Component 1 — No left-half rectangle with #004D4D fill found. "
                  f"Found {len(rectangles)} rectangle(s)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # =========================================================================
    # Component 2: Title 'Innovation Summit 2026' — 40pt bold Montserrat white (0.35 points)
    # =========================================================================
    try:
        title_found = False
        title_score = 0.0
        for shape in textboxes:
            full_text = shape.text_frame.text.strip()
            if "Innovation Summit 2026" in full_text:
                title_found = True
                # Check font properties on the runs containing the title text
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "Innovation Summit" in run.text:
                            # Check font name
                            if run.font.name and "Montserrat" in run.font.name:
                                title_score += 0.07
                                print(f"  PASS: Title font name = {run.font.name}")
                            else:
                                print(f"  FAIL: Title font name = {run.font.name}, expected Montserrat")

                            # Check font size (~40pt = 508000 EMU)
                            if run.font.size is not None:
                                expected_size = Pt(40)  # 508000
                                if abs(run.font.size - expected_size) <= Pt(1):
                                    title_score += 0.07
                                    print(f"  PASS: Title font size = {run.font.size} (~40pt)")
                                else:
                                    print(f"  FAIL: Title font size = {run.font.size}, expected {expected_size}")
                            else:
                                print(f"  FAIL: Title font size is None")

                            # Check bold
                            if run.font.bold is True:
                                title_score += 0.07
                                print(f"  PASS: Title is bold")
                            else:
                                print(f"  FAIL: Title bold = {run.font.bold}, expected True")

                            # Check color white (#FFFFFF)
                            try:
                                if run.font.color.type is not None:
                                    rgb_str = str(run.font.color.rgb).upper()
                                    if rgb_str == "FFFFFF":
                                        title_score += 0.07
                                        print(f"  PASS: Title color = #{rgb_str}")
                                    else:
                                        print(f"  FAIL: Title color = #{rgb_str}, expected #FFFFFF")
                                else:
                                    print(f"  FAIL: Title color type is None")
                            except Exception:
                                print(f"  FAIL: Title color not accessible")

                            # Check text content exactly
                            if full_text == "Innovation Summit 2026":
                                title_score += 0.07
                                print(f"  PASS: Title text exact match")
                            else:
                                print(f"  FAIL: Title text = '{full_text}', expected 'Innovation Summit 2026'")
                            break
                    if title_found:
                        break
                break

        if title_found:
            print(f"PASS: Component 2 — Title verification subtotal: {title_score:.2f}/0.35 pts")
            total_score += title_score
        else:
            print(f"FAIL: Component 2 — Title 'Innovation Summit 2026' not found in any text shape")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # =========================================================================
    # Component 3: Subtitle 'Shaping Tomorrow Together' — 22pt italic #004D4D (0.30 points)
    # =========================================================================
    try:
        subtitle_found = False
        subtitle_score = 0.0
        for shape in textboxes:
            full_text = shape.text_frame.text.strip()
            if "Shaping Tomorrow Together" in full_text:
                subtitle_found = True
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "Shaping Tomorrow" in run.text:
                            # Check font name
                            if run.font.name and "Montserrat" in run.font.name:
                                subtitle_score += 0.06
                                print(f"  PASS: Subtitle font name = {run.font.name}")
                            else:
                                print(f"  FAIL: Subtitle font name = {run.font.name}, expected Montserrat")

                            # Check font size (~22pt = 279400 EMU)
                            if run.font.size is not None:
                                expected_size = Pt(22)  # 279400
                                if abs(run.font.size - expected_size) <= Pt(1):
                                    subtitle_score += 0.06
                                    print(f"  PASS: Subtitle font size = {run.font.size} (~22pt)")
                                else:
                                    print(f"  FAIL: Subtitle font size = {run.font.size}, expected {expected_size}")
                            else:
                                print(f"  FAIL: Subtitle font size is None")

                            # Check italic
                            if run.font.italic is True:
                                subtitle_score += 0.06
                                print(f"  PASS: Subtitle is italic")
                            else:
                                print(f"  FAIL: Subtitle italic = {run.font.italic}, expected True")

                            # Check color #004D4D
                            try:
                                if run.font.color.type is not None:
                                    rgb_str = str(run.font.color.rgb).upper()
                                    if rgb_str == "004D4D":
                                        subtitle_score += 0.06
                                        print(f"  PASS: Subtitle color = #{rgb_str}")
                                    else:
                                        print(f"  FAIL: Subtitle color = #{rgb_str}, expected #004D4D")
                                else:
                                    print(f"  FAIL: Subtitle color type is None")
                            except Exception:
                                print(f"  FAIL: Subtitle color not accessible")

                            # Check text content exactly
                            if full_text == "Shaping Tomorrow Together":
                                subtitle_score += 0.06
                                print(f"  PASS: Subtitle text exact match")
                            else:
                                print(f"  FAIL: Subtitle text = '{full_text}', expected 'Shaping Tomorrow Together'")
                            break
                    if subtitle_found:
                        break
                break

        if subtitle_found:
            print(f"PASS: Component 3 — Subtitle verification subtotal: {subtitle_score:.2f}/0.30 pts")
            total_score += subtitle_score
        else:
            print(f"FAIL: Component 3 — Subtitle 'Shaping Tomorrow Together' not found in any text shape")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # =========================================================================
    # Component 4: Spatial layout — title on left half, subtitle on right half (0.10 points)
    # =========================================================================
    try:
        title_shape = None
        subtitle_shape = None
        for shape in textboxes:
            text = shape.text_frame.text.strip()
            if "Innovation Summit 2026" in text:
                title_shape = shape
            elif "Shaping Tomorrow Together" in text:
                subtitle_shape = shape

        layout_score = 0.0
        if title_shape is not None and subtitle_shape is not None:
            # Title center should be on the left half
            title_center_x = title_shape.left + title_shape.width / 2
            if title_center_x < midpoint_x:
                layout_score += 0.05
                print(f"  PASS: Title center X ({title_center_x}) is on left half (midpoint={midpoint_x})")
            else:
                print(f"  FAIL: Title center X ({title_center_x}) is NOT on left half (midpoint={midpoint_x})")

            # Subtitle center should be on the right half
            subtitle_center_x = subtitle_shape.left + subtitle_shape.width / 2
            if subtitle_center_x > midpoint_x:
                layout_score += 0.05
                print(f"  PASS: Subtitle center X ({subtitle_center_x}) is on right half")
            else:
                print(f"  FAIL: Subtitle center X ({subtitle_center_x}) is NOT on right half")

            if layout_score > 0:
                print(f"PASS: Component 4 — Spatial layout subtotal: {layout_score:.2f}/0.10 pts")
                total_score += layout_score
        else:
            print(f"FAIL: Component 4 — Cannot verify layout: "
                  f"title_shape={'found' if title_shape else 'missing'}, "
                  f"subtitle_shape={'found' if subtitle_shape else 'missing'}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
