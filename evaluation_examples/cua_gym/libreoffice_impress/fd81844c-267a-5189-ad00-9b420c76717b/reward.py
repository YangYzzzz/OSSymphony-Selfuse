"""
FINAL REWARD SCRIPT - SUCCESS
Task: I've got a presentation in LibreOffice Impress, and I want slide 2 to stand out with a solid light-blue color background. Can someone guide me on how to set it to #CFE8FF?
Generated: 2025-08-07 09:32:09
Status: success
Model: o4-mini
Total Steps: 2
"""

import os
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.dml.color import RGBColor

def verify_impress_task(file_path):
    print(f"Starting Impress verification for: {file_path}")
    total_score = 0.0
    max_score = 1.0

    # 1. File existence check (0.2 points)
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0
    print("✓ File exists (0.2 points)")
    total_score += 0.2

    # 2. Load presentation (0.1 points)
    try:
        presentation = Presentation(file_path)
        print(f"✓ Presentation loaded with {len(presentation.slides)} slides (0.1 points)")
        total_score += 0.1
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        print(f"REWARD: {min(total_score, max_score)}")
        return min(total_score, max_score)

    # 3. Slide count check (0.1 points)
    slide_count = len(presentation.slides)
    if slide_count >= 2:
        print("✓ Slide 2 exists (0.1 points)")
        total_score += 0.1
    else:
        print(f"✗ Slide count insufficient: found {slide_count} slides")

    # 4. Background fill type check (0.3 points)
    try:
        slide2 = presentation.slides[1]
        bg_fill = slide2.background.fill
        fill_type = getattr(bg_fill, 'type', None)
        if fill_type == MSO_FILL.SOLID:
            print("✓ Slide 2 background fill type is SOLID (0.3 points)")
            total_score += 0.3
        else:
            print(f"✗ Slide 2 background fill type is not SOLID: {fill_type}")
    except Exception as e:
        print(f"✗ Error checking fill type: {e}")

    # 5. Background color check (0.3 points)
    try:
        fc = bg_fill.fore_color
        rgb = getattr(fc, 'rgb', None)
        color_hex = None
        if isinstance(rgb, RGBColor):
            color_hex = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
        elif isinstance(rgb, str) and len(rgb) == 6:
            color_hex = f"#{rgb.upper()}"
        else:
            # Attempt to convert other types to string if unexpected
            rgb_str = str(rgb)
            if len(rgb_str) == 6 and all(c in '0123456789ABCDEFabcdef' for c in rgb_str):
                color_hex = f"#{rgb_str.upper()}"
        if color_hex == '#CFE8FF':
            print(f"✓ Slide 2 background color is correct ({color_hex}) (0.3 points)")
            total_score += 0.3
        else:
            print(f"✗ Slide 2 background color incorrect: {color_hex}")
    except Exception as e:
        print(f"✗ Error checking background color: {e}")

    # Final score calculation and output
    final_score = min(total_score, max_score)
    print(f"Final score: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score

if __name__ == '__main__':
    file_path = '/home/user/ive_got_a_presentation_in_libreoffice_impress_and_i_want_slide_2_to_stand_out_with_a_solid_light_blu.pptx'
    verify_impress_task(file_path)
