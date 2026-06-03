"""
Initial Setup: Create opening slide of ProductLaunch_X1 presentation
Task ID: impress_ps_012
Domain: libreoffice_impress

Creates a 10-slide presentation. Slide 1 is empty with white background.
Slides 2-10 contain realistic product launch content.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_012'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                 font_color=RGBColor(0x00, 0x00, 0x00), font_name="Calibri"):
    """Helper to add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = font_color
    return txBox


def create_initial():
    prs = Presentation()

    # Use blank layout (index 6) for all slides
    blank_layout = prs.slide_layouts[6]

    # --- Slide 1: Empty with white background (task target) ---
    slide1 = prs.slides.add_slide(blank_layout)
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # Slide 1 is intentionally empty - no text, no shapes

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(blank_layout)
    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    add_text_box(slide2, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
                 "Agenda", font_size=36, bold=True,
                 font_color=RGBColor(0x1A, 0x23, 0x7E), alignment=PP_ALIGN.LEFT)
    agenda_items = [
        "1. Product Overview & Key Features",
        "2. Technical Specifications",
        "3. Market Positioning & Competitive Analysis",
        "4. Pricing Strategy",
        "5. Go-to-Market Timeline",
        "6. Q&A"
    ]
    y_pos = Inches(1.6)
    for item in agenda_items:
        add_text_box(slide2, Inches(1.2), y_pos, Inches(7), Inches(0.4),
                     item, font_size=18, font_color=RGBColor(0x33, 0x33, 0x33))
        y_pos += Inches(0.5)

    # --- Slide 3: Product Overview ---
    slide3 = prs.slides.add_slide(blank_layout)
    add_text_box(slide3, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
                 "Product Overview", font_size=32, bold=True,
                 font_color=RGBColor(0x1A, 0x23, 0x7E), alignment=PP_ALIGN.LEFT)
    add_text_box(slide3, Inches(0.8), Inches(1.4), Inches(8), Inches(3.5),
                 "The X1 Pro represents our next-generation flagship device, "
                 "engineered from the ground up to deliver uncompromising performance. "
                 "Built on a 4nm process node with our proprietary NeuralCore architecture, "
                 "the X1 Pro pushes the boundaries of what's possible in mobile computing. "
                 "With 16GB of LPDDR5X RAM and up to 1TB of UFS 4.0 storage, "
                 "it handles demanding workloads with ease.",
                 font_size=16, font_color=RGBColor(0x44, 0x44, 0x44))

    # --- Slide 4: Key Features ---
    slide4 = prs.slides.add_slide(blank_layout)
    add_text_box(slide4, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
                 "Key Features", font_size=32, bold=True,
                 font_color=RGBColor(0x1A, 0x23, 0x7E))
    features = [
        "6.7-inch ProMotion LTPO OLED Display (2800x1260)",
        "NeuralCore X1 Chipset - 40% faster than predecessor",
        "108MP Triple Camera System with OIS",
        "5500mAh Battery with 120W HyperCharge",
        "IP68 Water & Dust Resistance",
        "Wi-Fi 7 and Bluetooth 5.4 Support"
    ]
    y_pos = Inches(1.5)
    for feat in features:
        add_text_box(slide4, Inches(1.2), y_pos, Inches(7.5), Inches(0.4),
                     feat, font_size=16, font_color=RGBColor(0x33, 0x33, 0x33))
        y_pos += Inches(0.5)

    # --- Slide 5: Technical Specifications ---
    slide5 = prs.slides.add_slide(blank_layout)
    add_text_box(slide5, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
                 "Technical Specifications", font_size=32, bold=True,
                 font_color=RGBColor(0x1A, 0x23, 0x7E))
    specs = [
        "Processor: NeuralCore X1, 4nm, 8-core (1x3.5GHz + 3x2.8GHz + 4x1.8GHz)",
        "Memory: 16GB LPDDR5X RAM",
        "Storage: 256GB / 512GB / 1TB UFS 4.0",
        "Display: 6.7\" LTPO OLED, 1-120Hz adaptive refresh",
        "Battery: 5500mAh, 120W wired / 50W wireless charging",
        "OS: Android 16 with X1 UI 5.0"
    ]
    y_pos = Inches(1.5)
    for spec in specs:
        add_text_box(slide5, Inches(1.0), y_pos, Inches(8), Inches(0.4),
                     spec, font_size=14, font_color=RGBColor(0x44, 0x44, 0x44))
        y_pos += Inches(0.5)

    # --- Slide 6: Market Analysis ---
    slide6 = prs.slides.add_slide(blank_layout)
    add_text_box(slide6, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
                 "Market Positioning", font_size=32, bold=True,
                 font_color=RGBColor(0x1A, 0x23, 0x7E))
    add_text_box(slide6, Inches(0.8), Inches(1.4), Inches(8), Inches(3.0),
                 "The premium smartphone market grew 12% YoY in 2025, reaching $198B globally. "
                 "The X1 Pro targets the $800-$1200 segment where Samsung Galaxy S26 Ultra "
                 "and iPhone 17 Pro Max currently dominate. Our differentiation centers on "
                 "computational photography capabilities and battery longevity.",
                 font_size=16, font_color=RGBColor(0x44, 0x44, 0x44))

    # --- Slide 7: Pricing Strategy ---
    slide7 = prs.slides.add_slide(blank_layout)
    add_text_box(slide7, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
                 "Pricing Strategy", font_size=32, bold=True,
                 font_color=RGBColor(0x1A, 0x23, 0x7E))
    prices = [
        "X1 Pro 256GB  -  $899",
        "X1 Pro 512GB  -  $999",
        "X1 Pro 1TB    -  $1,149",
        "",
        "Early bird discount: 15% off for pre-orders before Feb 28, 2026",
        "Trade-in program: Up to $400 credit for qualifying devices"
    ]
    y_pos = Inches(1.5)
    for price in prices:
        if price:
            add_text_box(slide7, Inches(1.2), y_pos, Inches(7), Inches(0.4),
                         price, font_size=18, font_color=RGBColor(0x33, 0x33, 0x33))
        y_pos += Inches(0.45)

    # --- Slide 8: Go-to-Market Timeline ---
    slide8 = prs.slides.add_slide(blank_layout)
    add_text_box(slide8, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
                 "Go-to-Market Timeline", font_size=32, bold=True,
                 font_color=RGBColor(0x1A, 0x23, 0x7E))
    timeline = [
        "Jan 15, 2026  -  Media embargo lift & press event",
        "Feb 1, 2026   -  Pre-orders open globally",
        "Feb 20, 2026  -  Reviewer units shipped",
        "Mar 1, 2026   -  Retail channel stocking begins",
        "Mar 15, 2026  -  Official launch day",
        "Apr 1, 2026   -  Expansion to 35 additional markets"
    ]
    y_pos = Inches(1.5)
    for item in timeline:
        add_text_box(slide8, Inches(1.0), y_pos, Inches(8), Inches(0.4),
                     item, font_size=16, font_color=RGBColor(0x44, 0x44, 0x44))
        y_pos += Inches(0.5)

    # --- Slide 9: Partnership & Distribution ---
    slide9 = prs.slides.add_slide(blank_layout)
    add_text_box(slide9, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
                 "Distribution Partners", font_size=32, bold=True,
                 font_color=RGBColor(0x1A, 0x23, 0x7E))
    partners = [
        "Carrier Partners: T-Mobile, Verizon, AT&T, Vodafone, Deutsche Telekom",
        "Retail: Best Buy, Amazon, Target, MediaMarkt",
        "Online: Direct-to-consumer via x1pro.com",
        "Enterprise: CDW, Insight, SHI International"
    ]
    y_pos = Inches(1.5)
    for partner in partners:
        add_text_box(slide9, Inches(1.0), y_pos, Inches(8), Inches(0.4),
                     partner, font_size=16, font_color=RGBColor(0x44, 0x44, 0x44))
        y_pos += Inches(0.55)

    # --- Slide 10: Q&A ---
    slide10 = prs.slides.add_slide(blank_layout)
    fill10 = slide10.background.fill
    fill10.solid()
    fill10.fore_color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    add_text_box(slide10, Inches(1.5), Inches(2.5), Inches(7), Inches(1.5),
                 "Questions & Answers", font_size=40, bold=True,
                 font_color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_text_box(slide10, Inches(2), Inches(4.2), Inches(6), Inches(0.6),
                 "product-team@company.com", font_size=20,
                 font_color=RGBColor(0xB0, 0xB0, 0xB0), alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
