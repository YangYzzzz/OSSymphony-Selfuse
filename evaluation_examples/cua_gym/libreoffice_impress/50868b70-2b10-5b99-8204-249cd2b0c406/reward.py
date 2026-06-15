"""
Reward Script: Product launch opening slide verification
Task ID: impress_ps_012
Domain: libreoffice_impress
Scoring:
  - Component 1: Gradient background (dark blue to black) — 0.25 pts
  - Component 2: 'X1 Pro' text (54pt, white, bold, centered) — 0.30 pts
  - Component 3: 'Redefining Performance' (28pt, light gray, italic, centered) — 0.25 pts
  - Component 4: 'Available March 15, 2026' (18pt, white, centered) — 0.20 pts
"""

import os
import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_012'


def get_all_text_shapes(slide):
    """Recursively get all text shapes including those in groups."""
    def extract(shape):
        results = []
        if hasattr(shape, "text") and hasattr(shape, "text_frame"):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.extend(extract(sub))
        return results
    out = []
    for shape in slide.shapes:
        out.extend(extract(shape))
    return out


def check_gradient_background(pptx_path):
    """Check if slide 1 has a gradient background with dark blue to black colors."""
    ns = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            with zf.open('ppt/slides/slide1.xml') as f:
                root = ET.parse(f).getroot()
                bg = root.find('.//p:bg', ns)
                if bg is None:
                    return False, "No background element found"
                grad = bg.find('.//a:gradFill', ns)
                if grad is None:
                    return False, "No gradient fill found on slide 1 background"
                gs_list = grad.findall('.//a:gs', ns)
                if len(gs_list) < 2:
                    return False, f"Gradient has {len(gs_list)} stops, expected at least 2"
                # Extract gradient stop colors
                colors = []
                for gs in gs_list:
                    clr = gs.find('.//a:srgbClr', ns)
                    if clr is not None:
                        colors.append(clr.get('val', '').upper())
                # Check for dark blue and black colors (with some tolerance)
                has_dark_blue = False
                has_black = False
                for c in colors:
                    if c in ('1A237E', '1A1A2E', '0D1B2A', '1B1B3A', '000033', '00004D', '0A0A3C', '0D0D4D'):
                        has_dark_blue = True
                    # Also accept approximate dark blues (R<0x30, G<0x30, B>0x50)
                    try:
                        r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                        if r < 0x40 and g < 0x40 and b > 0x50:
                            has_dark_blue = True
                        if r < 0x20 and g < 0x20 and b < 0x20:
                            has_black = True
                    except:
                        pass
                if has_dark_blue and has_black:
                    return True, f"Gradient with dark blue and black stops: {colors}"
                else:
                    return False, f"Gradient colors don't match: {colors}, need dark blue + black"
    except Exception as e:
        return False, f"Error checking gradient: {e}"


def find_text_shape(shapes, target_text):
    """Find shape containing text that matches target (case-insensitive, stripped)."""
    target_lower = target_text.strip().lower()
    for shape in shapes:
        if shape.has_text_frame:
            full_text = shape.text_frame.text.strip().lower()
            if full_text == target_lower:
                return shape
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    if len(prs.slides) == 0:
        print("CRITICAL: No slides in presentation")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]
    text_shapes = get_all_text_shapes(slide)

    # Component 1: Gradient background (dark blue to black) — 0.25 points
    try:
        passed, details = check_gradient_background(file_path)
        if passed:
            print(f"PASS: Component 1 — Gradient background verified: {details} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — {details}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: 'X1 Pro' text (54pt, white, bold, centered) — 0.30 points
    try:
        shape = find_text_shape(text_shapes, "X1 Pro")
        if shape is None:
            print("FAIL: Component 2 — 'X1 Pro' text not found on slide 1")
        else:
            comp2_score = 0.0
            para = shape.text_frame.paragraphs[0]
            runs = [r for r in para.runs if (r.text or "").strip()]

            if not runs:
                print("FAIL: Component 2 — No runs found in 'X1 Pro' text box")
            else:
                run = runs[0]
                # Check text content
                if "x1 pro" in run.text.strip().lower():
                    comp2_score += 0.06
                    print(f"  PASS: 'X1 Pro' text found: {repr(run.text)}")
                else:
                    print(f"  FAIL: Expected 'X1 Pro', found: {repr(run.text)}")

                # Check font size (54pt = 685800 EMU)
                if run.font.size is not None and abs(run.font.size - Pt(54)) < Pt(2):
                    comp2_score += 0.06
                    print(f"  PASS: Font size ~54pt (actual: {run.font.size / 12700}pt)")
                else:
                    actual_size = run.font.size / 12700 if run.font.size else None
                    print(f"  FAIL: Expected 54pt, found: {actual_size}pt")

                # Check white color
                try:
                    if run.font.color.type is not None and str(run.font.color.rgb).upper() == 'FFFFFF':
                        comp2_score += 0.06
                        print("  PASS: White font color")
                    else:
                        color_val = str(run.font.color.rgb) if run.font.color.type is not None else "None/inherited"
                        print(f"  FAIL: Expected white (FFFFFF), found: {color_val}")
                except Exception:
                    print("  FAIL: Could not read font color")

                # Check bold
                if run.font.bold is True:
                    comp2_score += 0.06
                    print("  PASS: Bold enabled")
                else:
                    print(f"  FAIL: Expected bold=True, found: {run.font.bold}")

                # Check centered alignment
                align = para.alignment
                if align == PP_ALIGN.CENTER:
                    comp2_score += 0.06
                    print("  PASS: Center aligned")
                else:
                    print(f"  FAIL: Expected CENTER, found: {align}")

            total_score += comp2_score
            print(f"  Component 2 subtotal: {comp2_score}/0.30")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: 'Redefining Performance' (28pt, light gray #B0B0B0, italic, centered) — 0.25 points
    try:
        shape = find_text_shape(text_shapes, "Redefining Performance")
        if shape is None:
            print("FAIL: Component 3 — 'Redefining Performance' text not found on slide 1")
        else:
            comp3_score = 0.0
            para = shape.text_frame.paragraphs[0]
            runs = [r for r in para.runs if (r.text or "").strip()]

            if not runs:
                print("FAIL: Component 3 — No runs found in 'Redefining Performance' text box")
            else:
                run = runs[0]
                # Check text content
                if "redefining performance" in run.text.strip().lower():
                    comp3_score += 0.05
                    print(f"  PASS: 'Redefining Performance' text found")
                else:
                    print(f"  FAIL: Expected 'Redefining Performance', found: {repr(run.text)}")

                # Check font size (28pt = 355600 EMU)
                if run.font.size is not None and abs(run.font.size - Pt(28)) < Pt(2):
                    comp3_score += 0.05
                    print(f"  PASS: Font size ~28pt (actual: {run.font.size / 12700}pt)")
                else:
                    actual_size = run.font.size / 12700 if run.font.size else None
                    print(f"  FAIL: Expected 28pt, found: {actual_size}pt")

                # Check light gray color (B0B0B0, with tolerance)
                try:
                    if run.font.color.type is not None:
                        rgb_str = str(run.font.color.rgb).upper()
                        r, g, b = int(rgb_str[0:2], 16), int(rgb_str[2:4], 16), int(rgb_str[4:6], 16)
                        # Accept light gray: R, G, B all in [0x90, 0xD0] range and close to each other
                        if (0x90 <= r <= 0xD0 and 0x90 <= g <= 0xD0 and 0x90 <= b <= 0xD0 and
                                abs(r - g) < 0x20 and abs(g - b) < 0x20):
                            comp3_score += 0.05
                            print(f"  PASS: Light gray font color ({rgb_str})")
                        else:
                            print(f"  FAIL: Expected light gray (~B0B0B0), found: {rgb_str}")
                    else:
                        print("  FAIL: Font color not set (inherited)")
                except Exception:
                    print("  FAIL: Could not read font color")

                # Check italic
                if run.font.italic is True:
                    comp3_score += 0.05
                    print("  PASS: Italic enabled")
                else:
                    print(f"  FAIL: Expected italic=True, found: {run.font.italic}")

                # Check centered alignment
                align = para.alignment
                if align == PP_ALIGN.CENTER:
                    comp3_score += 0.05
                    print("  PASS: Center aligned")
                else:
                    print(f"  FAIL: Expected CENTER, found: {align}")

            total_score += comp3_score
            print(f"  Component 3 subtotal: {comp3_score}/0.25")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: 'Available March 15, 2026' (18pt, white, centered) — 0.20 points
    try:
        shape = find_text_shape(text_shapes, "Available March 15, 2026")
        if shape is None:
            print("FAIL: Component 4 — 'Available March 15, 2026' text not found on slide 1")
        else:
            comp4_score = 0.0
            para = shape.text_frame.paragraphs[0]
            runs = [r for r in para.runs if (r.text or "").strip()]

            if not runs:
                print("FAIL: Component 4 — No runs found in release date text box")
            else:
                run = runs[0]
                # Check text content
                if "available march 15, 2026" in run.text.strip().lower():
                    comp4_score += 0.05
                    print(f"  PASS: Release date text found")
                else:
                    print(f"  FAIL: Expected 'Available March 15, 2026', found: {repr(run.text)}")

                # Check font size (18pt = 228600 EMU)
                if run.font.size is not None and abs(run.font.size - Pt(18)) < Pt(2):
                    comp4_score += 0.05
                    print(f"  PASS: Font size ~18pt (actual: {run.font.size / 12700}pt)")
                else:
                    actual_size = run.font.size / 12700 if run.font.size else None
                    print(f"  FAIL: Expected 18pt, found: {actual_size}pt")

                # Check white color
                try:
                    if run.font.color.type is not None and str(run.font.color.rgb).upper() == 'FFFFFF':
                        comp4_score += 0.05
                        print("  PASS: White font color")
                    else:
                        color_val = str(run.font.color.rgb) if run.font.color.type is not None else "None/inherited"
                        print(f"  FAIL: Expected white (FFFFFF), found: {color_val}")
                except Exception:
                    print("  FAIL: Could not read font color")

                # Check centered alignment
                align = para.alignment
                if align == PP_ALIGN.CENTER:
                    comp4_score += 0.05
                    print("  PASS: Center aligned")
                else:
                    print(f"  FAIL: Expected CENTER, found: {align}")

            total_score += comp4_score
            print(f"  Component 4 subtotal: {comp4_score}/0.20")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook: save any unsaved GUI edits before verification
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
