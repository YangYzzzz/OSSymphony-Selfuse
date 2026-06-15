"""
Initial Setup: Safety Training presentation with 15 slides, no progress bars.
Task ID: impress_ps_039
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_039'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                alignment=PP_ALIGN.LEFT, color=None, font_name="Arial"):
    """Helper to add a styled textbox."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14):
    """Add a bulleted text list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(font_size)
    return txBox


def create_initial():
    prs = Presentation()
    # Standard 16:9 slide dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Title Only layout
    full_blank = prs.slide_layouts[5]     # Blank layout

    slide_content = [
        {
            "title": "Workplace Safety Training 2025",
            "subtitle": "Comprehensive Safety Program\nAcme Manufacturing Inc.",
            "layout": "title"
        },
        {
            "title": "Training Objectives",
            "bullets": [
                "Understand OSHA workplace safety regulations",
                "Identify common hazards in manufacturing environments",
                "Learn proper use of personal protective equipment (PPE)",
                "Practice emergency evacuation procedures",
                "Report incidents using the new digital system"
            ]
        },
        {
            "title": "Safety Statistics — Q4 2024",
            "bullets": [
                "Total recordable incidents: 12 (down 23% from Q3)",
                "Lost-time injuries: 3 cases averaging 4.2 days each",
                "Near-miss reports filed: 87 (up 45% — good reporting culture)",
                "Safety audit compliance rate: 94.7%",
                "Workers' compensation costs: $148,320"
            ]
        },
        {
            "title": "Hazard Identification",
            "bullets": [
                "Chemical exposure: solvent fumes in painting bay B-7",
                "Slip/trip hazards: wet floors near wash stations",
                "Electrical risks: damaged cords on portable equipment",
                "Ergonomic issues: repetitive motion at assembly line 3",
                "Falling objects: overhead crane load path areas"
            ]
        },
        {
            "title": "Personal Protective Equipment (PPE)",
            "bullets": [
                "Hard hats required in all production zones",
                "Safety glasses with side shields — ANSI Z87.1 rated",
                "Steel-toed boots mandatory on shop floor",
                "Hearing protection in areas above 85 dB",
                "Chemical-resistant gloves for handling solvents"
            ]
        },
        {
            "title": "Machine Guarding & Lockout/Tagout",
            "bullets": [
                "All rotating equipment must have guards in place before startup",
                "LOTO procedure: Notify → Shut down → Isolate → Lock → Tag → Verify",
                "Only authorized personnel may remove locks",
                "Annual revalidation of LOTO procedures per OSHA 29 CFR 1910.147",
                "New interlock system installed on CNC machines in Hall D"
            ]
        },
        {
            "title": "Fire Safety & Prevention",
            "bullets": [
                "Fire extinguisher locations: every 75 feet along corridors",
                "Types: ABC dry chemical (general), CO2 (electrical), Class K (kitchen)",
                "Monthly inspection by safety team — check pin, gauge, condition",
                "Hot work permits required for welding and cutting operations",
                "No storage within 3 feet of electrical panels"
            ]
        },
        {
            "title": "Emergency Evacuation Plan",
            "bullets": [
                "Primary assembly point: Parking Lot C (north entrance)",
                "Secondary assembly point: Athletic field (east campus)",
                "Evacuation wardens assigned per department — see roster",
                "Mobility-impaired personnel: designated safe areas on each floor",
                "Quarterly evacuation drills — next drill: March 28, 2025"
            ]
        },
        {
            "title": "Incident Reporting Procedure",
            "bullets": [
                "Report ALL incidents within 2 hours — including near misses",
                "Use SafeTrack app or kiosk terminals at each entrance",
                "Supervisor notification required within 4 hours",
                "Investigation team deployed within 24 hours for recordable events",
                "Root cause analysis completed within 5 business days"
            ]
        },
        {
            "title": "Ergonomics & Wellness",
            "bullets": [
                "Workstation assessments available upon request (HR portal)",
                "Stretch breaks every 90 minutes for repetitive tasks",
                "Adjustable monitor arms installed in all office areas",
                "Anti-fatigue mats provided for standing workstations",
                "Employee assistance program (EAP): 1-800-555-SAFE"
            ]
        },
        {
            "title": "Chemical Safety & SDS Access",
            "bullets": [
                "Safety Data Sheets available at every chemical storage area",
                "Digital SDS library: intranet.acme.com/sds",
                "GHS label reading: pictograms, signal words, hazard statements",
                "Spill kit locations marked with yellow diamond signs",
                "Annual chemical inventory audit scheduled for April 2025"
            ]
        },
        {
            "title": "Confined Space Entry",
            "bullets": [
                "Permit required for all confined space entry — no exceptions",
                "Atmospheric testing: O2 (19.5-23.5%), LEL (<10%), H2S (<10 ppm)",
                "Trained attendant must remain at entry point at all times",
                "Rescue plan reviewed and signed before each entry",
                "Equipment: harness, tripod retrieval system, gas monitor"
            ]
        },
        {
            "title": "Electrical Safety",
            "bullets": [
                "Only qualified electricians may work on energized equipment",
                "Arc flash assessment labels on all panels > 50V",
                "GFCI protection required for all outdoor and wet location outlets",
                "Extension cords: temporary use only — max 90 days",
                "Report damaged cords and outlets immediately to Maintenance"
            ]
        },
        {
            "title": "Safety Culture & Recognition",
            "bullets": [
                "Monthly safety star award: $200 gift card + parking spot",
                "Department with lowest incident rate wins quarterly pizza lunch",
                "Anonymous safety suggestion box: feedback reviewed weekly",
                "New peer-to-peer safety observation program launching Q2",
                "Management safety walks: minimum twice per week per area"
            ]
        },
        {
            "title": "Summary & Next Steps",
            "bullets": [
                "Complete online quiz by March 15, 2025 (passing score: 80%)",
                "Schedule hands-on fire extinguisher training with your supervisor",
                "Update your emergency contact information in HR portal",
                "Review department-specific safety procedures with team lead",
                "Questions? Contact Safety Department: safety@acme.com"
            ]
        },
    ]

    for i, content in enumerate(slide_content):
        if i == 0 and content.get("layout") == "title":
            # Title slide
            slide = prs.slides.add_slide(full_blank)
            # Background color for title slide
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0x00, 0x3D, 0x6B)  # dark navy blue

            # Title
            add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(2),
                        content["title"], font_size=40, bold=True,
                        alignment=PP_ALIGN.CENTER, color=RGBColor(0xFF, 0xFF, 0xFF))

            # Subtitle
            txBox = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9), Inches(1.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            for j, line in enumerate(content["subtitle"].split("\n")):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.alignment = PP_ALIGN.CENTER
                run = p.runs[0]
                run.font.name = "Arial"
                run.font.size = Pt(22)
                run.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
        else:
            # Content slide
            slide = prs.slides.add_slide(full_blank)

            # Title bar area with accent color
            title_bg = slide.shapes.add_shape(
                1,  # MSO_SHAPE.RECTANGLE
                Inches(0), Inches(0), prs.slide_width, Inches(1.3)
            )
            title_fill = title_bg.fill
            title_fill.solid()
            title_fill.fore_color.rgb = RGBColor(0x00, 0x3D, 0x6B)
            title_bg.line.fill.background()

            add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(1),
                        content["title"], font_size=28, bold=True,
                        alignment=PP_ALIGN.LEFT, color=RGBColor(0xFF, 0xFF, 0xFF))

            # Slide number in corner
            add_textbox(slide, Inches(12), Inches(0.3), Inches(1), Inches(0.5),
                        str(i + 1), font_size=14, bold=False,
                        alignment=PP_ALIGN.RIGHT, color=RGBColor(0xCC, 0xCC, 0xCC))

            # Bullet content
            if "bullets" in content:
                add_bullet_list(slide, Inches(1), Inches(1.8), Inches(11), Inches(4.5),
                                content["bullets"], font_size=16)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Number of slides: {len(prs.slides)}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
