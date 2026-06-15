"""
Reward Script: Progress bar on every slide in Safety_Training.pptx
Task ID: impress_ps_039
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): All 15 slides have a full-width gray background bar at the bottom
  Component 2 (0.4): All 15 slides have a colored progress bar with width proportional to slide_number/15
  Component 3 (0.3): Progress bar colors and positioning are correct (green fill, proper top/height)
"""

import os

from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_039'
TOTAL_SLIDES = 15


def find_bottom_rectangles(slide, slide_width):
    """Find rectangle shapes near the bottom of the slide that form the progress bar."""
    gray_bar = None
    colored_bar = None

    for shape in slide.shapes:
        # Look for AUTO_SHAPE rectangles near the bottom (top >= 6500000 EMU)
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or shape.shape_type == 1:
            if shape.top >= 6500000:
                try:
                    fill = shape.fill
                    if fill.type is not None and fill.type == 1:  # SOLID fill
                        rgb = str(fill.fore_color.rgb).upper()
                        # Full-width gray bar
                        if rgb in ('D9D9D9', 'C0C0C0', 'BFBFBF', 'E0E0E0', 'D3D3D3', 'CCCCCC'):
                            if shape.width >= slide_width * 0.95:
                                gray_bar = shape
                        # Colored progress bar (green-ish or accent)
                        else:
                            colored_bar = shape
                except Exception:
                    pass
    return gray_bar, colored_bar


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    slide_width = prs.slide_width

    # Precondition: must have 15 slides
    if num_slides != TOTAL_SLIDES:
        print(f"FAIL: Expected {TOTAL_SLIDES} slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: All 15 slides have a full-width gray background bar at the bottom (0.3 points)
    try:
        gray_bar_count = 0
        for i, slide in enumerate(prs.slides):
            gray_bar, _ = find_bottom_rectangles(slide, slide_width)
            if gray_bar is not None:
                gray_bar_count += 1
            else:
                print(f"  FAIL: Slide {i+1} missing gray background bar")

        if gray_bar_count > 0:
            comp1_score = 0.3 * (gray_bar_count / TOTAL_SLIDES)
            if gray_bar_count == TOTAL_SLIDES:
                print(f"PASS: Component 1 -- All {TOTAL_SLIDES} slides have gray background bar ({comp1_score:.2f} pts)")
            else:
                print(f"PARTIAL: Component 1 -- {gray_bar_count}/{TOTAL_SLIDES} slides have gray background bar ({comp1_score:.2f} pts)")
            total_score += comp1_score
        else:
            print(f"FAIL: Component 1 -- No slides have gray background bar (0.00 pts)")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: All 15 slides have colored progress bar with correct proportional width (0.4 points)
    try:
        width_correct_count = 0
        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            _, colored_bar = find_bottom_rectangles(slide, slide_width)
            if colored_bar is None:
                print(f"  FAIL: Slide {slide_num} missing colored progress bar")
                continue

            # Expected width: slide_num / TOTAL_SLIDES * slide_width
            expected_width = int(slide_num * slide_width / TOTAL_SLIDES)
            actual_width = colored_bar.width
            # Allow 5% tolerance
            tolerance = 0.05
            if expected_width == 0:
                width_ok = actual_width == 0
            else:
                width_ok = abs(actual_width - expected_width) / expected_width <= tolerance

            if width_ok:
                width_correct_count += 1
            else:
                print(f"  FAIL: Slide {slide_num} progress bar width={actual_width}, expected ~{expected_width} (ratio={actual_width/slide_width:.3f}, expected={slide_num/TOTAL_SLIDES:.3f})")

        if width_correct_count > 0:
            comp2_score = 0.4 * (width_correct_count / TOTAL_SLIDES)
            if width_correct_count == TOTAL_SLIDES:
                print(f"PASS: Component 2 -- All {TOTAL_SLIDES} slides have correct progress bar width ({comp2_score:.2f} pts)")
            else:
                print(f"PARTIAL: Component 2 -- {width_correct_count}/{TOTAL_SLIDES} slides have correct width ({comp2_score:.2f} pts)")
            total_score += comp2_score
        else:
            print(f"FAIL: Component 2 -- No slides have correct progress bar width (0.00 pts)")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Progress bar positioning and color correctness (0.3 points)
    # Check: colored bars have a green-ish fill, bars are at consistent vertical position, proper height
    try:
        position_correct_count = 0
        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            gray_bar, colored_bar = find_bottom_rectangles(slide, slide_width)

            if gray_bar is None or colored_bar is None:
                print(f"  FAIL: Slide {slide_num} missing bar(s) for position/color check")
                continue

            checks_ok = True

            # Check colored bar has a solid fill color (not gray)
            try:
                fill = colored_bar.fill
                if fill.type != 1:
                    checks_ok = False
                    print(f"  FAIL: Slide {slide_num} colored bar fill type is not solid")
            except Exception:
                checks_ok = False

            # Check bars are at same vertical position
            if gray_bar.top != colored_bar.top:
                checks_ok = False
                print(f"  FAIL: Slide {slide_num} bars at different top positions ({gray_bar.top} vs {colored_bar.top})")

            # Check bar height is reasonable (thin bar, < 200000 EMU = ~0.55 cm)
            if colored_bar.height > 200000 or colored_bar.height < 30000:
                checks_ok = False
                print(f"  FAIL: Slide {slide_num} bar height unusual: {colored_bar.height}")

            # Check colored bar starts at left=0
            if colored_bar.left != 0:
                checks_ok = False
                print(f"  FAIL: Slide {slide_num} colored bar left={colored_bar.left}, expected 0")

            if checks_ok:
                position_correct_count += 1

        if position_correct_count > 0:
            comp3_score = 0.3 * (position_correct_count / TOTAL_SLIDES)
            if position_correct_count == TOTAL_SLIDES:
                print(f"PASS: Component 3 -- All {TOTAL_SLIDES} slides have correct bar positioning/color ({comp3_score:.2f} pts)")
            else:
                print(f"PARTIAL: Component 3 -- {position_correct_count}/{TOTAL_SLIDES} slides correct ({comp3_score:.2f} pts)")
            total_score += comp3_score
        else:
            print(f"FAIL: Component 3 -- No slides have correct positioning/color (0.00 pts)")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    # Try alternate name
    file_path = f'{WORKDIR}/Safety_Training.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {WORKDIR}/{TASK_ID}.pptx or {WORKDIR}/Safety_Training.pptx")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
