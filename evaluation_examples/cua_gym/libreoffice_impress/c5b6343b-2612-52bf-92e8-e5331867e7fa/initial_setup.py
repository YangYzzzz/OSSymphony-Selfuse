"""
Initial Setup: Create a 12-slide presentation for PDF export task
Task ID: impstruct_039
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
OUTPUT = f'{WORKDIR}/full_presentation.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q1 2025 Business Review"
    slide1.placeholders[1].text = "Prepared by Strategic Planning Division\nMarch 2025"

    # Slide 2: Executive Summary
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Revenue grew 18% year-over-year to $142.3M"
    p2a = tf2.add_paragraph()
    p2a.text = "Customer acquisition cost decreased by 12%"
    p2b = tf2.add_paragraph()
    p2b.text = "Net promoter score improved from 62 to 71"
    p2c = tf2.add_paragraph()
    p2c.text = "Three new product lines launched successfully"

    # Slide 3: Revenue Breakdown
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Revenue Breakdown by Region"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "North America: $68.4M (48%)"
    for line in [
        "Europe: $39.8M (28%)",
        "Asia-Pacific: $24.2M (17%)",
        "Latin America: $9.9M (7%)",
    ]:
        p = tf3.add_paragraph()
        p.text = line

    # Slide 4: Customer Metrics
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Customer Metrics Dashboard"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Total Active Customers: 14,832 (+22% YoY)"
    for line in [
        "Enterprise Tier: 342 accounts ($78.2M ARR)",
        "Mid-Market: 1,245 accounts ($41.6M ARR)",
        "SMB: 13,245 accounts ($22.5M ARR)",
        "Churn Rate: 3.2% (down from 4.8%)",
    ]:
        p = tf4.add_paragraph()
        p.text = line

    # Slide 5: Product Updates
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Product Development Updates"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Platform v3.2 released with AI-powered analytics"
    for line in [
        "Mobile app redesign completed (4.7 star rating)",
        "API gateway performance improved by 340%",
        "Security audit passed with zero critical findings",
    ]:
        p = tf5.add_paragraph()
        p.text = line

    # Slide 6: Team Performance
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Team Performance Highlights"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Engineering: 94% sprint velocity achievement"
    for line in [
        "Sales: 112% of quarterly target reached",
        "Support: Average response time reduced to 1.4 hours",
        "Marketing: 2.1M website visitors (+35% QoQ)",
    ]:
        p = tf6.add_paragraph()
        p.text = line

    # Slide 7: Financial Overview
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Financial Overview"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "Operating Expenses: $98.7M"
    for line in [
        "EBITDA: $43.6M (30.6% margin)",
        "Free Cash Flow: $31.2M",
        "R&D Investment: $28.4M (19.9% of revenue)",
    ]:
        p = tf7.add_paragraph()
        p.text = line

    # Slide 8: Market Analysis
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Competitive Landscape"
    tf8 = slide8.placeholders[1].text_frame
    tf8.text = "Market share increased to 23.4% (from 19.8%)"
    for line in [
        "Competitor A: 31.2% (stable)",
        "Competitor B: 18.7% (declining)",
        "New entrant C captured 4.1% market share",
    ]:
        p = tf8.add_paragraph()
        p.text = line

    # Slide 9: Strategic Initiatives
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Q2 Strategic Initiatives"
    tf9 = slide9.placeholders[1].text_frame
    tf9.text = "Launch enterprise data warehouse integration"
    for line in [
        "Expand APAC sales team by 15 headcount",
        "Complete SOC 2 Type II certification",
        "Begin Series D fundraising preparation",
    ]:
        p = tf9.add_paragraph()
        p.text = line

    # Slide 10: Risk Assessment
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "Risk Assessment"
    tf10 = slide10.placeholders[1].text_frame
    tf10.text = "Supply chain delays impacting hardware rollout"
    for line in [
        "Regulatory changes in EU data privacy laws",
        "Key talent retention in competitive market",
        "Currency fluctuation exposure in APAC region",
    ]:
        p = tf10.add_paragraph()
        p.text = line

    # Slide 11: Timeline & Milestones
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    slide11.shapes.title.text = "Key Milestones & Timeline"
    tf11 = slide11.placeholders[1].text_frame
    tf11.text = "April: Enterprise integration beta launch"
    for line in [
        "May: APAC office opening (Singapore)",
        "June: Annual customer conference (CloudConnect 2025)",
        "July: Platform v4.0 development kickoff",
    ]:
        p = tf11.add_paragraph()
        p.text = line

    # Slide 12: Thank You / Q&A
    slide12 = prs.slides.add_slide(prs.slide_layouts[0])
    slide12.shapes.title.text = "Thank You"
    slide12.placeholders[1].text = "Questions & Discussion\nContact: strategy@acmecorp.com"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=3.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
