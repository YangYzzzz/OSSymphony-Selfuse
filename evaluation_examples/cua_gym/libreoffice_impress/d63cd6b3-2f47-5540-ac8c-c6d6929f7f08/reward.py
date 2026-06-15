"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 129 the bullets are sitting way too close to the words. In LibreOffice Impress, how do I adjust the gap between the bullet symbol and its text to exactly 0.2 cm?
Generated: 2025-09-10 22:50:26
Status: success
Model: azure-o3
Total Steps: 7
"""

import os
from pptx import Presentation

# -------------------------------------------------------------
# Utility: convert centimetres to English Metric Units (EMU)
# 1 inch  = 2.54 cm ; 1 inch = 914400 EMU
# -------------------------------------------------------------

def cm_to_emu(cm: float) -> int:
    """Convert centimetres to EMUs (rounded to nearest integer)."""
    return int(round((cm / 2.54) * 914400))

# -------------------------------------------------------------
# Core verification logic
# -------------------------------------------------------------

def verify_bullet_gap(file_path: str,
                      slide_number: int = 129,
                      expected_gap_cm: float = 0.2,
                      tolerance_cm: float = 0.02) -> float:
    """Verify that the gap between a bullet symbol and the accompanying text
    (i.e., the hanging indent) on the specified slide equals the expected
    distance within a given tolerance.

    A progressive score between 0.0 and 1.0 is returned based on the share of
    bullet paragraphs that meet the requirement. A score of **exactly 1.0**
    is returned only when every detected bullet paragraph is correct.
    """

    print(f"Verifying bullet gap on slide {slide_number} → expected {expected_gap_cm} cm (±{tolerance_cm} cm)…")

    # ------------------------------------------------------------------
    # PRECONDITIONS  (No points awarded for these — they must simply pass)
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ Presentation file not found.")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        return 0.0

    if len(prs.slides) < slide_number:
        print(f"✗ Slide {slide_number} does not exist (presentation has only {len(prs.slides)} slides).")
        return 0.0

    slide = prs.slides[slide_number - 1]  # zero-based index

    expected_emu   = cm_to_emu(expected_gap_cm)
    tolerance_emu  = cm_to_emu(tolerance_cm)
    print(f"Expected hanging indent: {expected_emu} EMU (±{tolerance_emu})")

    # ------------------------------------------------------------------
    # Analyse every paragraph on the slide and measure bullet indents
    # ------------------------------------------------------------------
    bullet_paragraphs  = 0  # paragraphs that are (very likely) bullet items
    correct_paragraphs = 0  # bullet paragraphs that match the requirement

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for para in shape.text_frame.paragraphs:
            pPr = para._p.pPr
            if pPr is None:
                continue

            indent_attr = pPr.get('indent')
            if indent_attr is None:
                continue  # no indent attribute → unlikely to be a bullet

            # LibreOffice writes hanging indents as NEGATIVE values (in EMU)
            try:
                indent_val = int(indent_attr)
            except ValueError:
                continue  # malformed value → ignore

            if indent_val >= 0:
                continue  # positive or zero ⇒ not a hanging indent ⇒ skip

            # At this point we have a paragraph with a negative indent — treat
            # it as a bullet item for verification purposes
            bullet_paragraphs += 1
            gap_val = abs(indent_val)

            if abs(gap_val - expected_emu) <= tolerance_emu:
                correct_paragraphs += 1
            else:
                print(f"  ✗ Bullet indent {indent_val} EMU outside expected range.")

    # ------------------------------------------------------------------
    # Scoring
    # ------------------------------------------------------------------
    if bullet_paragraphs == 0:
        print("✗ No bullet paragraphs with a hanging indent detected on the slide.")
        return 0.0

    ratio = correct_paragraphs / bullet_paragraphs
    print(f"✓ Correct bullet gaps: {correct_paragraphs}/{bullet_paragraphs} → ratio = {ratio:.2f}")

    score = round(min(ratio, 1.0), 2)
    print(f"REWARD: {score}")
    return score

# -------------------------------------------------------------
# Execute verification (entry-point)
# -------------------------------------------------------------

if __name__ == "__main__":
    # Path provided by the task context
    FILE_PATH = "/home/user/on_slide_129_the_bullets_are_sitting_way_too_close_to_the_words_in_libreoffice_impress_how_do_i_adju_golden.pptx"
    verify_bullet_gap(FILE_PATH)

