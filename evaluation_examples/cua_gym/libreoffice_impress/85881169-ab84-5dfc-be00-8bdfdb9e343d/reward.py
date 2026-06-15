"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 6, at the end of the equation, I need to include the Greek letter omega (Ω). Could you guide me on how to insert this special character in LibreOffice Impress?
Generated: 2025-08-07 10:52:07
Status: success
Model: o4-mini
Total Steps: 1
"""

import os
from pptx import Presentation

def verify_impress_omega(file_path):
    print("Checking task completion: insert Omega on slide 6 equation.")
    score = 0.0
    max_score = 1.0

    # Requirement 1: File exists (0.2 points)
    if os.path.exists(file_path):
        print(f"✓ File exists: {file_path} (0.2 points)")
        score += 0.2
    else:
        print(f"✗ File not found: {file_path} (0 points)")
        print(f"REWARD: {score}")
        return score

    # Requirement 2: Load presentation (no direct points, preparatory)
    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded presentation, slide count: {len(prs.slides)}")
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        print(f"REWARD: {score}")
        return score

    # Requirement 3: Slide count >= 6 (0.2 points)
    if len(prs.slides) >= 6:
        print("✓ Presentation has at least 6 slides (0.2 points)")
        score += 0.2
        slide = prs.slides[5]
    else:
        print(f"✗ Presentation has fewer than 6 slides ({len(prs.slides)}) (0 points)")
        print(f"REWARD: {score}")
        return score

    # Requirement 4: Search for Omega character in text shapes on slide 6
    found_text = False
    omega_found = False
    print("Searching for shapes with text on slide 6...")

    for shape in slide.shapes:
        if not hasattr(shape, 'text'):
            continue
        text = shape.text.strip()
        if not text:
            continue
        found_text = True
        print(f"  Found text: '{text}'")
        if 'Ω' in text:
            omega_found = True
            print("  ✓ Omega character found in text (0.3 points)")
            score += 0.3
            if text.endswith('Ω'):
                print("  ✓ Omega is at the end of the text (0.3 points)")
                score += 0.3
            else:
                print("  ✗ Omega is not at the end of the text (0 points)")
            break

    if not found_text:
        print("✗ No text shapes found on slide 6 (0 points)")
    elif not omega_found:
        print("✗ No Omega character found in any text shape (0 points)")

    # Final score calculation
    final_score = min(score, max_score)
    print(f"Final score (capped at 1.0): {final_score}")
    return final_score

# Execute verification
def main():
    file_path = '/home/user/on_slide_6_at_the_end_of_the_equation_i_need_to_include_the_greek_letter_omega_ω_could_you_guide_me_.pptx'
    reward = verify_impress_omega(file_path)
    print(f"REWARD: {reward}")

if __name__ == '__main__':
    main()
