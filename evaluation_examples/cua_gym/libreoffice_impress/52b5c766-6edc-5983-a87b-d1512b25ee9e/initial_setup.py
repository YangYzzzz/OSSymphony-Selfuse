"""
Initial Setup: Create a 10-slide data presentation for impress_stu_078
Task ID: impress_stu_078
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_078'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, text in enumerate(bullets):
        if i == 0:
            body.paragraphs[0].text = text
        else:
            p = body.add_paragraph()
            p.text = text
            p.level = 0
    return slide


def add_blank_text_slide(prs, title_text, body_text, left=Inches(0.8), top=Inches(1.8)):
    """Add a slide with a title and a free-form text box."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title text box
    ttl = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = ttl.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # Body text box
    txBox = slide.shapes.add_textbox(left, top, Inches(8.4), Inches(4.5))
    tf2 = txBox.text_frame
    tf2.word_wrap = True
    lines = body_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            tf2.paragraphs[0].text = line
            tf2.paragraphs[0].space_after = Pt(6)
            for r in tf2.paragraphs[0].runs:
                r.font.size = Pt(16)
        else:
            p = tf2.add_paragraph()
            p.text = line
            p.space_after = Pt(6)
            for r in p.runs:
                r.font.size = Pt(16)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs,
                    "Q4 2025 Revenue Analysis",
                    "Prepared by Analytics Division — Confidential")

    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "Total Q4 revenue reached $14.8M, up 12% YoY",
        "North America contributed 58% of total revenue",
        "SaaS subscription revenue grew 23% quarter-over-quarter",
        "Customer acquisition cost decreased by 8%",
        "Net promoter score improved from 42 to 51",
    ])

    # Slide 3: Revenue Breakdown by Region
    add_content_slide(prs, "Revenue Breakdown by Region", [
        "North America: $8.58M (58%) — up 15% YoY",
        "Europe: $3.70M (25%) — up 9% YoY",
        "Asia-Pacific: $1.78M (12%) — up 18% YoY",
        "Latin America: $0.74M (5%) — up 6% YoY",
    ])

    # Slide 4: Monthly Trend Analysis
    add_content_slide(prs, "Monthly Revenue Trend — Q4 2025", [
        "October: $4.52M — Driven by annual renewal cycle",
        "November: $4.88M — Black Friday promotions boosted sales",
        "December: $5.40M — Year-end enterprise deals closed",
        "Average monthly growth rate: 9.3%",
    ])

    # Slide 5: Product Line Performance
    add_content_slide(prs, "Product Line Performance", [
        "Enterprise Suite: $6.2M (42%) — 340 active contracts",
        "Professional Plan: $4.1M (28%) — 1,240 subscribers",
        "Starter Plan: $2.8M (19%) — 5,680 subscribers",
        "Add-ons & Services: $1.7M (11%) — consulting + training",
    ])

    # Slide 6: Customer Metrics
    add_content_slide(prs, "Customer Acquisition & Retention", [
        "New customers acquired: 892 (up from 764 in Q3)",
        "Customer churn rate: 3.2% (down from 4.1%)",
        "Average contract value: $18,400",
        "Customer lifetime value: $62,300",
        "Upsell conversion rate: 28%",
    ])

    # Slide 7: Cost Analysis
    add_content_slide(prs, "Operational Cost Analysis", [
        "Total operating expenses: $9.8M",
        "R&D investment: $3.4M (35% of OpEx)",
        "Sales & Marketing: $3.1M (32% of OpEx)",
        "General & Administrative: $2.0M (20% of OpEx)",
        "Infrastructure: $1.3M (13% of OpEx)",
    ])

    # Slide 8: Competitive Landscape
    add_content_slide(prs, "Competitive Landscape", [
        "Market share increased from 12.4% to 13.8%",
        "Key competitor TechVault lost 2 enterprise accounts to us",
        "Feature parity achieved on AI-driven analytics module",
        "Pricing remains 15% below industry average",
    ])

    # Slide 9: Chart Analysis — THIS IS THE TARGET SLIDE
    # Use blank layout so we have full control of content placement
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    # Title
    ttl9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf_title = ttl9.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "Chart Analysis: Revenue vs. Expenses"
    p_title.alignment = PP_ALIGN.LEFT
    r_title = p_title.runs[0]
    r_title.font.size = Pt(28)
    r_title.font.bold = True
    r_title.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # Chart analysis body text
    body9 = slide9.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(3.5))
    tf9 = body9.text_frame
    tf9.word_wrap = True

    analysis_lines = [
        "The revenue-to-expense ratio improved from 1.32 in Q3 to 1.51 in Q4, indicating stronger operational efficiency.",
        "Revenue growth outpaced expense growth by 7.2 percentage points, primarily driven by SaaS margin improvements.",
        "Infrastructure costs remained flat despite a 23% increase in active users, validating our cloud migration strategy.",
        "The largest expense reduction came from automated customer onboarding, saving approximately $420K in Q4.",
        "Projected Q1 2026 margin expansion: 2-3% based on current trajectory and planned headcount optimization.",
    ]
    for i, line in enumerate(analysis_lines):
        if i == 0:
            tf9.paragraphs[0].text = line
            tf9.paragraphs[0].space_after = Pt(8)
            for r in tf9.paragraphs[0].runs:
                r.font.size = Pt(16)
        else:
            p = tf9.add_paragraph()
            p.text = line
            p.space_after = Pt(8)
            for r in p.runs:
                r.font.size = Pt(16)

    # NO hyperlink here — that's the task for the agent

    # Slide 10: Next Steps
    add_content_slide(prs, "Next Steps & Recommendations", [
        "Expand enterprise sales team by 4 account executives in Q1",
        "Launch AI analytics module beta to top 50 customers",
        "Renegotiate cloud infrastructure contract — potential 12% savings",
        "Pilot referral program targeting professional plan users",
        "Schedule quarterly business review with top 20 accounts",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
