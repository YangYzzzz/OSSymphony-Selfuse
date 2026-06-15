"""
Reward Script: Insert hyperlink on slide 9
Task ID: impress_stu_078
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30) - Hyperlink text "Link to full dataset" exists on slide 9
  Component 2 (0.30) - Hyperlink URL is correct (https://docs.google.com/spreadsheets/example)
  Component 3 (0.20) - Font size is 14pt (177800 EMU)
  Component 4 (0.20) - Font is blue (0000FF) and underlined
"""

import os
from pptx import Presentation
from pptx.util import Pt

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_078'


def persist_app_state(domain):
    """Save any unsaved LibreOffice state before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def find_hyperlink_shapes(slide):
    """Find all shapes on a slide that contain a hyperlink in any run."""
    results = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.hyperlink and run.hyperlink.address:
                    results.append({
                        'shape': shape,
                        'para': para,
                        'run': run,
                        'text': run.text,
                        'url': run.hyperlink.address,
                    })
    return results


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 9 slides
    if len(prs.slides) < 9:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 9")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[8]  # slide 9 (0-indexed)

    # Find hyperlink shapes on slide 9
    hl_shapes = find_hyperlink_shapes(slide)

    # Component 1: Hyperlink with display text "Link to full dataset" exists on slide 9 (0.30 points)
    target_hl = None
    try:
        for hl in hl_shapes:
            if 'link to full dataset' in hl['text'].lower().strip():
                target_hl = hl
                break

        if target_hl is not None:
            print(f"PASS: Component 1 - Found hyperlink text '{target_hl['text']}' on slide 9 (0.30 pts)")
            total_score += 0.30
        else:
            # Also check all text on slide 9 for partial match (text exists but no hyperlink)
            all_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        all_texts.append(para.text)
            print(f"FAIL: Component 1 - No hyperlink with text 'Link to full dataset' found on slide 9")
            print(f"  Available texts: {all_texts}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Hyperlink URL is correct (0.30 points)
    try:
        if target_hl is not None:
            expected_url = 'https://docs.google.com/spreadsheets/example'
            actual_url = target_hl['url']
            if actual_url.rstrip('/') == expected_url.rstrip('/'):
                print(f"PASS: Component 2 - Hyperlink URL is correct: {actual_url} (0.30 pts)")
                total_score += 0.30
            else:
                print(f"FAIL: Component 2 - Expected URL '{expected_url}', found '{actual_url}'")
        else:
            print(f"FAIL: Component 2 - No target hyperlink found, cannot check URL")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Font size is 14pt (177800 EMU) (0.20 points)
    try:
        if target_hl is not None:
            run = target_hl['run']
            actual_size = run.font.size
            expected_size = Pt(14)  # 177800 EMU
            if actual_size is not None and abs(actual_size - expected_size) <= 12700:  # 1pt tolerance
                print(f"PASS: Component 3 - Font size is {actual_size} EMU (~{actual_size/12700:.1f}pt), expected 14pt (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 3 - Expected font size ~14pt (177800 EMU), found {actual_size}")
        else:
            print(f"FAIL: Component 3 - No target hyperlink found, cannot check font size")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Font is blue and underlined (0.20 points)
    try:
        if target_hl is not None:
            run = target_hl['run']
            # Check underline
            is_underlined = run.font.underline is True
            # Check blue color
            is_blue = False
            try:
                if run.font.color.type is not None:
                    rgb = str(run.font.color.rgb).upper()
                    # Accept various shades of blue (0000FF, 0000EE, etc.)
                    is_blue = rgb == '0000FF' or (rgb[0:4] == '0000' and int(rgb[4:6], 16) >= 0xCC)
            except Exception:
                pass

            if is_underlined and is_blue:
                print(f"PASS: Component 4 - Font is blue and underlined (0.20 pts)")
                total_score += 0.20
            elif is_underlined:
                print(f"PARTIAL: Component 4 - Font is underlined but not blue (0.10 pts)")
                total_score += 0.10
            elif is_blue:
                print(f"PARTIAL: Component 4 - Font is blue but not underlined (0.10 pts)")
                total_score += 0.10
            else:
                color_info = "unknown"
                try:
                    color_info = str(run.font.color.rgb) if run.font.color.type is not None else "None/inherited"
                except:
                    pass
                print(f"FAIL: Component 4 - Expected blue underlined text, found underline={run.font.underline}, color={color_info}")
        else:
            print(f"FAIL: Component 4 - No target hyperlink found, cannot check font style")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state('libreoffice_impress')

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
