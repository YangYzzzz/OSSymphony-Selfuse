"""
FINAL REWARD SCRIPT - SUCCESS
Task: I’m tidying up slide 279 of my LibreOffice Impress deck, and the three pictures on that slide don’t line up nicely—one’s towering while another is stubby. What’s the quickest way to set every image’s height to precisely 5.0 cm so they’re uniform?
Generated: 2025-09-10 20:23:57
Status: success
Model: azure-o3
Total Steps: 1
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os


def verify_images_height(pptx_path: str, target_cm: float = 5.0, tolerance_cm: float = 0.15) -> float:
    """Reward-script verifier

    Checks slide 279 (0-based index 278) of the given PPTX to ensure **all** picture
    shapes have a height of `target_cm` ± `tolerance_cm` centimetres.

    Progressive scoring (max 1.0):
      • 0.2  – slide 279 exists
      • 0.3  – at least three pictures found on that slide
      • 0.5  – every picture height within tolerance of the target
    The function prints detailed diagnostics and returns the final score.
    """

    MAX_SCORE = 1.0
    score = 0.0

    print(f"Checking presentation file: {pptx_path}")

    # --- prerequisite checks (no points) -----------------------------------
    if not os.path.exists(pptx_path):
        print("✗ File does not exist – task failed.")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"✗ Error loading PPTX: {e}")
        print("REWARD: 0.0")
        return 0.0

    total_slides = len(prs.slides)
    print(f"Total slides in deck: {total_slides}")

    # --- requirement 1: slide 279 exists -----------------------------------
    target_index = 278  # 0-based for slide 279
    if total_slides <= target_index:
        print("✗ Slide 279 not found – task failed.")
        print("REWARD: 0.0")
        return 0.0
    print("✓ Slide 279 exists (0.2 points)")
    score += 0.2

    slide = prs.slides[target_index]

    # --- requirement 2: at least three pictures on the slide --------------
    pictures = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    print(f"Found {len(pictures)} picture shapes on slide 279.")

    if len(pictures) >= 3:
        print("✓ At least three pictures present (0.3 points)")
        score += 0.3
    else:
        print("✗ Fewer than three pictures – no points for this criterion")

    # --- requirement 3: every picture height == 5.0 cm ± tolerance ---------
    emu_per_cm = 360_000  # EMUs in one centimetre
    target_emu = target_cm * emu_per_cm
    tolerance_emu = tolerance_cm * emu_per_cm

    all_within_tolerance = True
    for idx, pic in enumerate(pictures, 1):
        height_cm = pic.height / emu_per_cm
        diff_emu = abs(pic.height - target_emu)
        within = diff_emu <= tolerance_emu
        status = "✓" if within else "✗"
        print(f"  {status} Picture {idx}: height = {height_cm:.2f} cm (Δ {diff_emu} EMU)")
        if not within:
            all_within_tolerance = False

    if pictures and all_within_tolerance:
        print("✓ All picture heights within tolerance (0.5 points)")
        score += 0.5
    elif pictures:
        print("✗ Not all pictures meet the 5.0 cm height requirement")

    # --- final score -------------------------------------------------------
    final_score = min(score, MAX_SCORE)
    print(f"Total score: {final_score}/{MAX_SCORE}")
    print(f"REWARD: {final_score}")
    return final_score


# ---------------------------------------------------------------------------
# Execute verification when run as a script
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    PRESENTATION_PATH = (
        "/home/user/"
        "im_tidying_up_slide_279_of_my_libreoffice_impress_deck_"
        "and_the_three_pictures_on_that_slide_dont_lin_golden.pptx"
    )
    verify_images_height(PRESENTATION_PATH)

