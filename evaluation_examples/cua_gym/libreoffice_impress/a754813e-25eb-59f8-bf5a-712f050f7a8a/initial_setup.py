"""
Initial Setup: Create Corporate_Standard.otp template with footer '2024 Confidential' and red accent color #CC0000
Task ID: impress_fix_043
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_043'
TEMPLATES_DIR = f'{WORKDIR}/Templates'
OUTPUT_OTP = f'{TEMPLATES_DIR}/Corporate_Standard.otp'
TEMP_PPTX = f'/tmp/{TASK_ID}_template.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn

    os.makedirs(TEMPLATES_DIR, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    RED_ACCENT = RGBColor(0xCC, 0x00, 0x00)
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

    # ─── Slide 1: Title Slide ───
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Corporate Strategy 2024"
    for run in slide1.shapes.title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = DARK_GRAY
        run.font.name = "Arial"
    slide1.placeholders[1].text = "Annual Planning & Review"
    for run in slide1.placeholders[1].text_frame.paragraphs[0].runs:
        run.font.size = Pt(18)
        run.font.color.rgb = RED_ACCENT
        run.font.name = "Arial"

    # Red accent bar at top of slide 1
    bar1 = slide1.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0), Inches(0), prs.slide_width, Inches(0.3)
    )
    bar1.fill.solid()
    bar1.fill.fore_color.rgb = RED_ACCENT
    bar1.line.fill.background()
    bar1.name = "AccentBar_Top"

    # Footer textbox on slide 1
    footer1 = slide1.shapes.add_textbox(
        Inches(0.5), Inches(6.8), Inches(4), Inches(0.5)
    )
    tf = footer1.text_frame
    p = tf.paragraphs[0]
    p.text = "2024 Confidential"
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = DARK_GRAY
        run.font.name = "Arial"
    footer1.name = "Footer_Text"

    # ─── Slide 2: Agenda ───
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    # Title
    title2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(1))
    tf2 = title2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Agenda"
    p2.alignment = PP_ALIGN.LEFT
    for run in p2.runs:
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = DARK_GRAY
        run.font.name = "Arial"

    # Red accent line under title
    line2 = slide2.shapes.add_shape(
        1, Inches(0.8), Inches(1.3), Inches(3), Inches(0.05)
    )
    line2.fill.solid()
    line2.fill.fore_color.rgb = RED_ACCENT
    line2.line.fill.background()
    line2.name = "AccentLine_Agenda"

    # Agenda items
    agenda_items = [
        "1. Market Analysis & Competitive Landscape",
        "2. Revenue Growth Projections for Q1-Q4",
        "3. Product Roadmap & Feature Prioritization",
        "4. Team Expansion & Talent Acquisition",
        "5. Budget Allocation & Resource Planning",
        "6. Risk Assessment & Mitigation Strategies",
    ]
    agenda_box = slide2.shapes.add_textbox(
        Inches(1.0), Inches(1.6), Inches(7), Inches(4)
    )
    atf = agenda_box.text_frame
    atf.word_wrap = True
    for i, item in enumerate(agenda_items):
        if i == 0:
            p = atf.paragraphs[0]
        else:
            p = atf.add_paragraph()
        p.text = item
        p.space_after = Pt(12)
        for run in p.runs:
            run.font.size = Pt(16)
            run.font.color.rgb = DARK_GRAY
            run.font.name = "Arial"

    # Footer on slide 2
    footer2 = slide2.shapes.add_textbox(
        Inches(0.5), Inches(6.8), Inches(4), Inches(0.5)
    )
    fp2 = footer2.text_frame.paragraphs[0]
    fp2.text = "2024 Confidential"
    fp2.alignment = PP_ALIGN.LEFT
    for run in fp2.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = DARK_GRAY
        run.font.name = "Arial"
    footer2.name = "Footer_Text"

    # ─── Slide 3: Financial Overview ───
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])

    title3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(1))
    tf3 = title3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Financial Overview"
    p3.alignment = PP_ALIGN.LEFT
    for run in p3.runs:
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = DARK_GRAY
        run.font.name = "Arial"

    # Accent line
    line3 = slide3.shapes.add_shape(
        1, Inches(0.8), Inches(1.3), Inches(3), Inches(0.05)
    )
    line3.fill.solid()
    line3.fill.fore_color.rgb = RED_ACCENT
    line3.line.fill.background()
    line3.name = "AccentLine_Financial"

    # Table with financial data
    table_shape = slide3.shapes.add_table(
        6, 4, Inches(0.8), Inches(1.6), Inches(8), Inches(3)
    )
    table = table_shape.table
    headers = ["Department", "Q1 Revenue", "Q2 Revenue", "Growth %"]
    data = [
        ["Engineering Services", "$2,340,000", "$2,780,000", "+18.8%"],
        ["Cloud Solutions", "$1,890,000", "$2,150,000", "+13.8%"],
        ["Consulting", "$945,000", "$1,120,000", "+18.5%"],
        ["Support & Maintenance", "$670,000", "$710,000", "+6.0%"],
        ["Total", "$5,845,000", "$6,760,000", "+15.7%"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(12)
            run.font.color.rgb = WHITE
            run.font.name = "Arial"
        # Red header background
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = RED_ACCENT

    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(11)
                run.font.color.rgb = DARK_GRAY
                run.font.name = "Arial"
                if r == len(data):
                    run.font.bold = True

    # Footer on slide 3
    footer3 = slide3.shapes.add_textbox(
        Inches(0.5), Inches(6.8), Inches(4), Inches(0.5)
    )
    fp3 = footer3.text_frame.paragraphs[0]
    fp3.text = "2024 Confidential"
    fp3.alignment = PP_ALIGN.LEFT
    for run in fp3.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = DARK_GRAY
        run.font.name = "Arial"
    footer3.name = "Footer_Text"

    # ─── Slide 4: Key Metrics ───
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])

    title4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(1))
    tf4 = title4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Key Performance Metrics"
    p4.alignment = PP_ALIGN.LEFT
    for run in p4.runs:
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = DARK_GRAY
        run.font.name = "Arial"

    # Accent line
    line4 = slide4.shapes.add_shape(
        1, Inches(0.8), Inches(1.3), Inches(3), Inches(0.05)
    )
    line4.fill.solid()
    line4.fill.fore_color.rgb = RED_ACCENT
    line4.line.fill.background()
    line4.name = "AccentLine_Metrics"

    # Metric cards (red-bordered boxes)
    metrics = [
        ("Customer Satisfaction", "94.2%", "+3.1%"),
        ("Employee Retention", "91.7%", "+2.4%"),
        ("Project Delivery", "88.5%", "+5.2%"),
    ]
    for i, (label, value, change) in enumerate(metrics):
        x = Inches(0.8 + i * 3.0)
        y = Inches(2.0)

        # Card background
        card = slide4.shapes.add_shape(1, x, y, Inches(2.5), Inches(2.0))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = RED_ACCENT
        card.line.width = Pt(2)
        card.name = f"MetricCard_{i}"

        # Metric value
        val_box = slide4.shapes.add_textbox(x, Inches(2.3), Inches(2.5), Inches(0.8))
        vp = val_box.text_frame.paragraphs[0]
        vp.text = value
        vp.alignment = PP_ALIGN.CENTER
        for run in vp.runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RED_ACCENT
            run.font.name = "Arial"

        # Metric label
        lbl_box = slide4.shapes.add_textbox(x, Inches(3.0), Inches(2.5), Inches(0.5))
        lp = lbl_box.text_frame.paragraphs[0]
        lp.text = label
        lp.alignment = PP_ALIGN.CENTER
        for run in lp.runs:
            run.font.size = Pt(12)
            run.font.color.rgb = DARK_GRAY
            run.font.name = "Arial"

        # Change indicator
        chg_box = slide4.shapes.add_textbox(x, Inches(3.4), Inches(2.5), Inches(0.4))
        cp = chg_box.text_frame.paragraphs[0]
        cp.text = change
        cp.alignment = PP_ALIGN.CENTER
        for run in cp.runs:
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0x00, 0x80, 0x00)
            run.font.name = "Arial"

    # Footer on slide 4
    footer4 = slide4.shapes.add_textbox(
        Inches(0.5), Inches(6.8), Inches(4), Inches(0.5)
    )
    fp4 = footer4.text_frame.paragraphs[0]
    fp4.text = "2024 Confidential"
    fp4.alignment = PP_ALIGN.LEFT
    for run in fp4.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = DARK_GRAY
        run.font.name = "Arial"
    footer4.name = "Footer_Text"

    # ─── Slide 5: Next Steps ───
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])

    title5 = slide5.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(1))
    tf5 = title5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Next Steps & Action Items"
    p5.alignment = PP_ALIGN.LEFT
    for run in p5.runs:
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = DARK_GRAY
        run.font.name = "Arial"

    # Accent line
    line5 = slide5.shapes.add_shape(
        1, Inches(0.8), Inches(1.3), Inches(3), Inches(0.05)
    )
    line5.fill.solid()
    line5.fill.fore_color.rgb = RED_ACCENT
    line5.line.fill.background()
    line5.name = "AccentLine_NextSteps"

    # Action items
    actions = [
        "Finalize Q3 budget proposals by March 28, 2024",
        "Complete hiring for three senior engineering positions",
        "Launch beta testing for CloudSync platform upgrade",
        "Schedule quarterly review with executive leadership",
        "Distribute updated compliance training materials",
    ]
    actions_box = slide5.shapes.add_textbox(
        Inches(1.0), Inches(1.6), Inches(7), Inches(4)
    )
    actf = actions_box.text_frame
    actf.word_wrap = True
    for i, item in enumerate(actions):
        if i == 0:
            p = actf.paragraphs[0]
        else:
            p = actf.add_paragraph()
        p.text = item
        p.space_after = Pt(14)
        for run in p.runs:
            run.font.size = Pt(14)
            run.font.color.rgb = DARK_GRAY
            run.font.name = "Arial"

    # Red accent bar at bottom
    bar_bottom = slide5.shapes.add_shape(
        1, Inches(0), Inches(7.2), prs.slide_width, Inches(0.3)
    )
    bar_bottom.fill.solid()
    bar_bottom.fill.fore_color.rgb = RED_ACCENT
    bar_bottom.line.fill.background()
    bar_bottom.name = "AccentBar_Bottom"

    # Footer on slide 5
    footer5 = slide5.shapes.add_textbox(
        Inches(0.5), Inches(6.8), Inches(4), Inches(0.5)
    )
    fp5 = footer5.text_frame.paragraphs[0]
    fp5.text = "2024 Confidential"
    fp5.alignment = PP_ALIGN.LEFT
    for run in fp5.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = DARK_GRAY
        run.font.name = "Arial"
    footer5.name = "Footer_Text"

    # Save as .pptx first, then convert to .otp
    prs.save(TEMP_PPTX)
    print(f'Temporary PPTX created: {TEMP_PPTX}')

    # Convert to OTP using LibreOffice
    result = subprocess.run(
        ['libreoffice', '--headless', '--convert-to', 'otp',
         '--outdir', TEMPLATES_DIR, TEMP_PPTX],
        capture_output=True, text=True, timeout=30,
        env={**os.environ, 'DISPLAY': ':0'}
    )
    print(f'Conversion stdout: {result.stdout}')
    print(f'Conversion stderr: {result.stderr}')

    # Rename the converted file to expected name
    converted_path = os.path.join(TEMPLATES_DIR, f'{TASK_ID}_template.otp')
    if os.path.exists(converted_path):
        os.rename(converted_path, OUTPUT_OTP)
        print(f'Renamed {converted_path} -> {OUTPUT_OTP}')
    elif os.path.exists(OUTPUT_OTP):
        print(f'Output already at expected path: {OUTPUT_OTP}')
    else:
        # Fallback: list what was produced
        for f in os.listdir(TEMPLATES_DIR):
            print(f'  Found: {TEMPLATES_DIR}/{f}')
        # Try alternative: just copy pptx as otp (LO can open it)
        import shutil
        shutil.copy(TEMP_PPTX, OUTPUT_OTP)
        print(f'Fallback: copied PPTX as OTP: {OUTPUT_OTP}')

    # Verify
    if os.path.exists(OUTPUT_OTP):
        sz = os.path.getsize(OUTPUT_OTP)
        print(f'Template created: {OUTPUT_OTP} ({sz} bytes)')
    else:
        print('ERROR: Template file was not created!')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT_OTP}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
