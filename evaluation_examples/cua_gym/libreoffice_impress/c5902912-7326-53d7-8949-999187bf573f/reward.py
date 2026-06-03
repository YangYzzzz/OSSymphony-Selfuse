"""
Reward Script: Product Launch Countdown Presentation
Task ID: impress_wf_078
Domain: libreoffice_impress
Scoring:
  C1 (0.10) - File exists on Desktop as Launch_Countdown.pptx
  C2 (0.10) - Exactly 8 slides
  C3 (0.20) - Slide 1: title text, large '30' text, red background (#F44336)
  C4 (0.10) - Slide 1: Zoom entrance animation on '30' text
  C5 (0.10) - Slide 2: horizontal arrow + 4 milestone markers
  C6 (0.10) - Slide 3: marketing checklist table (Task, Status, Owner)
  C7 (0.10) - Slide 4: 4 channel cards (Email, Social, PR, Events)
  C8 (0.10) - Slide 6: launch day agenda table (Time, Activity, Responsible)
  C9 (0.10) - Slide 7: metric placeholders and KPI cards
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_078'
FILE_PATH = os.path.join(WORKDIR, 'Desktop', 'Launch_Countdown.pptx')


def get_all_text_shapes(slide):
    """Recursively get all shapes with text, including inside groups."""
    def extract(shape):
        results = []
        if hasattr(shape, "text") and hasattr(shape, "text_frame"):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.extend(extract(sub))
        return results
    out = []
    for shape in slide.shapes:
        out.extend(extract(shape))
    return out


def check_animation_zoom(pptx_path, slide_idx=0):
    """Check if slide has a Zoom entrance animation (presetID=53, presetClass=entr)."""
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            slide_xml = f'ppt/slides/slide{slide_idx + 1}.xml'
            with zf.open(slide_xml) as f:
                root = ET.parse(f).getroot()
                # Find all cTn elements with presetID and presetClass
                for elem in root.iter(f'{{{ns_p}}}cTn'):
                    preset_id = elem.get('presetID')
                    preset_class = elem.get('presetClass')
                    if preset_id == '53' and preset_class == 'entr':
                        return True
        return False
    except Exception:
        return False


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    total_score = 0.0

    # Component 1: File exists on Desktop (0.10 points)
    # This differentiates initial (no file on Desktop) from golden
    try:
        if os.path.exists(file_path):
            print(f"PASS: Component 1 - File found at {file_path} (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 1 - File not found at {file_path}")
            print("REWARD: 0.0")
            return 0.0
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 2: Exactly 8 slides (0.10 points)
    try:
        num_slides = len(prs.slides)
        if num_slides == 8:
            print(f"PASS: Component 2 - Exactly 8 slides (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 - Expected 8 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Slide 1 - title text, large '30' text, red background (0.20 points)
    try:
        if num_slides >= 1:
            slide1 = prs.slides[0]
            sub_score = 0.0

            # 3a: Title contains "Product Launch" and "30 Days"
            all_text = [s.text for s in get_all_text_shapes(slide1)]
            has_title = any("Product Launch" in t and "30" in t for t in all_text)
            if has_title:
                sub_score += 0.07
                print(f"  PASS: C3a - Title text with 'Product Launch' and '30' found")
            else:
                print(f"  FAIL: C3a - Title text missing. Found texts: {all_text[:3]}")

            # 3b: Large '30' text element
            has_30 = any(t.strip() == "30" for t in all_text)
            if has_30:
                sub_score += 0.06
                print(f"  PASS: C3b - Large '30' text element found")
            else:
                print(f"  FAIL: C3b - No standalone '30' text found")

            # 3c: Red background (#F44336)
            bg = slide1.background.fill
            bg_color = None
            if bg.type == 1:  # SOLID
                bg_color = str(bg.fore_color.rgb)
            if bg_color and bg_color.upper() == "F44336":
                sub_score += 0.07
                print(f"  PASS: C3c - Red background #F44336 found")
            else:
                print(f"  FAIL: C3c - Expected background #F44336, found type={bg.type}, color={bg_color}")

            total_score += sub_score
            print(f"PASS: Component 3 - Slide 1 checks ({sub_score:.2f}/0.20 pts)")
        else:
            print(f"FAIL: Component 3 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Slide 1 Zoom animation on '30' text (0.10 points)
    try:
        has_zoom = check_animation_zoom(file_path, slide_idx=0)
        if has_zoom:
            print(f"PASS: Component 4 - Zoom entrance animation found on slide 1 (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 4 - No Zoom entrance animation on slide 1")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Slide 2 - horizontal arrow + 4 milestone markers (0.10 points)
    try:
        if num_slides >= 2:
            slide2 = prs.slides[1]
            has_arrow = False
            oval_count = 0
            for shape in slide2.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    name_lower = shape.name.lower()
                    if 'arrow' in name_lower:
                        has_arrow = True
                    if 'oval' in name_lower or 'circle' in name_lower:
                        oval_count += 1

            sub_score = 0.0
            if has_arrow:
                sub_score += 0.05
                print(f"  PASS: C5a - Arrow shape found on slide 2")
            else:
                print(f"  FAIL: C5a - No arrow shape on slide 2")

            if oval_count >= 4:
                sub_score += 0.05
                print(f"  PASS: C5b - {oval_count} milestone marker shapes found (>= 4)")
            else:
                print(f"  FAIL: C5b - Only {oval_count} milestone markers, expected >= 4")

            total_score += sub_score
            print(f"PASS: Component 5 - Slide 2 timeline ({sub_score:.2f}/0.10 pts)")
        else:
            print(f"FAIL: Component 5 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    # Component 6: Slide 3 - marketing checklist table (Task, Status, Owner) (0.10 points)
    try:
        if num_slides >= 3:
            slide3 = prs.slides[2]
            found_table = False
            for shape in slide3.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    if len(table.columns) >= 3 and len(table.rows) >= 2:
                        headers = [table.cell(0, c).text.strip().lower() for c in range(len(table.columns))]
                        has_task = any('task' in h for h in headers)
                        has_status = any('status' in h for h in headers)
                        has_owner = any('owner' in h for h in headers)
                        if has_task and has_status and has_owner:
                            found_table = True
                            break

            if found_table:
                print(f"PASS: Component 6 - Marketing checklist table found with Task/Status/Owner (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 6 - No valid marketing checklist table on slide 3")
        else:
            print(f"FAIL: Component 6 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 6 - {e}")

    # Component 7: Slide 4 - 4 channel cards (Email, Social, PR, Events) (0.10 points)
    try:
        if num_slides >= 4:
            slide4 = prs.slides[3]
            found_channels = set()
            expected_channels = {'email', 'social', 'pr', 'events'}
            for shape in slide4.shapes:
                if hasattr(shape, 'text') and shape.text.strip().lower() in expected_channels:
                    found_channels.add(shape.text.strip().lower())

            if found_channels >= expected_channels:
                print(f"PASS: Component 7 - All 4 channel cards found: {found_channels} (0.10 pts)")
                total_score += 0.10
            else:
                missing = expected_channels - found_channels
                print(f"FAIL: Component 7 - Missing channels: {missing}. Found: {found_channels}")
        else:
            print(f"FAIL: Component 7 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 7 - {e}")

    # Component 8: Slide 6 - Launch day agenda table (Time, Activity, Responsible) (0.10 points)
    try:
        if num_slides >= 6:
            slide6 = prs.slides[5]
            found_table = False
            for shape in slide6.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    if len(table.columns) >= 3 and len(table.rows) >= 2:
                        headers = [table.cell(0, c).text.strip().lower() for c in range(len(table.columns))]
                        has_time = any('time' in h for h in headers)
                        has_activity = any('activity' in h for h in headers)
                        has_responsible = any('responsible' in h for h in headers)
                        if has_time and has_activity and has_responsible:
                            found_table = True
                            break

            if found_table:
                print(f"PASS: Component 8 - Launch day agenda table found with Time/Activity/Responsible (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 8 - No valid agenda table on slide 6")
        else:
            print(f"FAIL: Component 8 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 8 - {e}")

    # Component 9: Slide 7 - metric placeholders and KPI cards (0.10 points)
    try:
        if num_slides >= 7:
            slide7 = prs.slides[6]
            kpi_count = 0
            placeholder_count = 0
            for shape in slide7.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    shape_text = shape.text.strip() if hasattr(shape, 'text') else ''
                    name_lower = shape.name.lower()
                    if 'rounded' in name_lower:
                        kpi_count += 1
                    elif 'rectangle' in name_lower or 'placeholder' in shape_text.lower() or 'chart' in shape_text.lower():
                        placeholder_count += 1

            sub_score = 0.0
            if kpi_count >= 3:
                sub_score += 0.05
                print(f"  PASS: C9a - {kpi_count} KPI card shapes found (>= 3)")
            else:
                print(f"  FAIL: C9a - Only {kpi_count} KPI cards, expected >= 3")

            if placeholder_count >= 1:
                sub_score += 0.05
                print(f"  PASS: C9b - {placeholder_count} chart placeholder shapes found (>= 1)")
            else:
                print(f"  FAIL: C9b - No chart placeholder shapes found")

            total_score += sub_score
            print(f"PASS: Component 9 - Slide 7 metrics ({sub_score:.2f}/0.10 pts)")
        else:
            print(f"FAIL: Component 9 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 9 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook: save any unsaved LibreOffice state before verification
def persist_app_state(domain):
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
persist_app_state("libreoffice_impress")

if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(FILE_PATH)
