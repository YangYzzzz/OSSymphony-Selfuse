"""
Initial Setup: Create tasks.json and blank presentation for weekly status report macro task.
Task ID: impress_gf5_032
Domain: libreoffice_impress
"""

import json
import os
import shlex
import subprocess
import time

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_032'
OUTPUT_PPTX = f'{WORKDIR}/{TASK_ID}.pptx'
TASKS_JSON = f'{WORKDIR}/tasks.json'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    # 1. Create tasks.json with 8 realistic tasks
    tasks = [
        {
            "task_name": "Migrate payment gateway to Stripe v3",
            "owner": "Sarah Chen",
            "status": "in_progress",
            "completion_pct": 75
        },
        {
            "task_name": "Redesign customer onboarding flow",
            "owner": "Marcus Johnson",
            "status": "completed",
            "completion_pct": 100
        },
        {
            "task_name": "Implement SSO for enterprise clients",
            "owner": "Priya Sharma",
            "status": "in_progress",
            "completion_pct": 40
        },
        {
            "task_name": "Fix memory leak in notification service",
            "owner": "David Kim",
            "status": "blocked",
            "completion_pct": 20
        },
        {
            "task_name": "Write API documentation for v2 endpoints",
            "owner": "Elena Rodriguez",
            "status": "in_progress",
            "completion_pct": 60
        },
        {
            "task_name": "Set up CI/CD pipeline for mobile app",
            "owner": "James O'Brien",
            "status": "completed",
            "completion_pct": 100
        },
        {
            "task_name": "Resolve database connection pooling issues",
            "owner": "Aisha Patel",
            "status": "blocked",
            "completion_pct": 10
        },
        {
            "task_name": "Deploy monitoring dashboards for prod",
            "owner": "Carlos Mendez",
            "status": "in_progress",
            "completion_pct": 55
        }
    ]

    with open(TASKS_JSON, 'w') as f:
        json.dump(tasks, f, indent=2)
    print(f'Created tasks.json with {len(tasks)} tasks: {TASKS_JSON}')

    # 2. Create a blank presentation (initial state before macro runs)
    prs = Presentation()
    # Just a single blank slide so Impress opens with something
    prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    prs.save(OUTPUT_PPTX)
    print(f'Created blank presentation: {OUTPUT_PPTX}')

    # 3. Open in LibreOffice Impress (GUI-ready state)
    launch_gui(f'libreoffice --impress "{OUTPUT_PPTX}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
