"""
Reward Script: Weekly Status Report Slide Deck from tasks.json
Task ID: impress_gf5_032
Domain: libreoffice_impress
Scoring:
  Component 1 (0.15): weekly_status.pptx exists and has correct slide count (10 slides)
  Component 2 (0.25): Slide 1 overview table with 9 rows x 4 columns, correct headers and all 8 tasks
  Component 3 (0.30): Slides 2-9 each have task name, owner/status info, and progress bar rectangles
  Component 4 (0.15): Progress bar widths are proportional to completion_pct
  Component 5 (0.15): Final slide (10) is a Blockers slide listing blocked tasks
"""

import os
import json

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_032'


def verify_task():
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0
    file_path = os.path.join(WORKDIR, 'weekly_status.pptx')

    # Precondition: file must exist
    if not os.path.exists(file_path):
        print(f"CRITICAL: weekly_status.pptx not found at {file_path}")
        print("REWARD: 0.0")
        return 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Load tasks.json for reference data
    try:
        with open(os.path.join(WORKDIR, 'tasks.json'), 'r') as f:
            tasks = json.load(f)
    except Exception as e:
        print(f"CRITICAL: Cannot load tasks.json: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_tasks = len(tasks)  # expected: 8
    blocked_tasks = [t for t in tasks if t['status'] == 'blocked']
    # Expected slides: 1 overview + 8 per-task + 1 blockers = 10
    expected_slides = 1 + num_tasks + 1

    # Component 1: File exists and has correct slide count (0.15 pts)
    try:
        num_slides = len(prs.slides)
        if num_slides == expected_slides:
            print(f"PASS: Component 1 — weekly_status.pptx has {num_slides} slides (expected {expected_slides}) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — expected {expected_slides} slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 1 overview table with all 8 tasks, 4 columns (0.25 pts)
    try:
        slide1 = prs.slides[0]
        # Find table shape
        table_shape = None
        for shape in slide1.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_shape = shape
                break

        if table_shape is None:
            print("FAIL: Component 2 — No table found on slide 1")
        else:
            table = table_shape.table
            num_rows = len(table.rows)
            num_cols = len(table.columns)

            # Check table dimensions: 9 rows (1 header + 8 tasks), 4 columns
            if num_rows == num_tasks + 1 and num_cols == 4:
                total_score += 0.10
                print(f"PASS: Component 2a — Table is {num_rows}x{num_cols} (correct) (0.10 pts)")
            else:
                print(f"FAIL: Component 2a — Table is {num_rows}x{num_cols}, expected {num_tasks+1}x4")

            # Check headers
            headers = [table.cell(0, c).text.strip().lower() for c in range(min(num_cols, 4))]
            expected_headers_keywords = ['task', 'owner', 'status', 'completion']
            headers_match = all(
                any(kw in h for h in headers)
                for kw in expected_headers_keywords
            )
            if headers_match:
                total_score += 0.05
                print(f"PASS: Component 2b — Headers contain expected keywords: {headers} (0.05 pts)")
            else:
                print(f"FAIL: Component 2b — Headers {headers} missing expected keywords")

            # Check all task names appear in the table
            table_texts = []
            for r in range(1, min(num_rows, num_tasks + 1)):
                row_text = ' '.join(table.cell(r, c).text for c in range(min(num_cols, 4)))
                table_texts.append(row_text.lower())

            tasks_found = 0
            for task in tasks:
                task_name_lower = task['task_name'].lower()
                if any(task_name_lower in rt for rt in table_texts):
                    tasks_found += 1

            if tasks_found == num_tasks:
                total_score += 0.10
                print(f"PASS: Component 2c — All {num_tasks} tasks found in table (0.10 pts)")
            else:
                print(f"FAIL: Component 2c — Only {tasks_found}/{num_tasks} tasks found in table")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slides 2-9 each have task name, owner/status, and progress bar rectangles (0.30 pts)
    try:
        per_slide_pts = 0.30 / num_tasks  # 0.0375 per slide

        for idx in range(num_tasks):
            slide_num = idx + 2  # slides 2-9 (0-indexed: 1-8)
            if slide_num - 1 >= len(prs.slides):
                print(f"FAIL: Component 3 — Slide {slide_num} does not exist")
                continue

            slide = prs.slides[slide_num - 1]
            task = tasks[idx]
            task_name = task['task_name']
            owner = task['owner']

            # Gather all text from slide
            all_text = ''
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        all_text += para.text + '\n'

            all_text_lower = all_text.lower()

            # Check task name present
            has_name = task_name.lower() in all_text_lower
            # Check owner present
            has_owner = owner.lower() in all_text_lower
            # Check progress bar rectangles (at least 2 AUTO_SHAPE)
            rect_count = sum(1 for s in slide.shapes if 'AUTO_SHAPE' in str(s.shape_type))
            has_bars = rect_count >= 2

            slide_ok = has_name and has_owner and has_bars
            if slide_ok:
                total_score += per_slide_pts
                print(f"PASS: Component 3 — Slide {slide_num} ({task_name[:30]}): name={has_name}, owner={has_owner}, bars={rect_count}")
            else:
                print(f"FAIL: Component 3 — Slide {slide_num} ({task_name[:30]}): name={has_name}, owner={has_owner}, bars={rect_count}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Progress bar widths proportional to completion_pct (0.15 pts)
    try:
        per_slide_pts4 = 0.15 / num_tasks
        tolerance = 0.05  # 5% tolerance

        for idx in range(num_tasks):
            slide_num = idx + 2
            if slide_num - 1 >= len(prs.slides):
                continue

            slide = prs.slides[slide_num - 1]
            task = tasks[idx]
            expected_pct = task['completion_pct'] / 100.0

            # Find the two rectangles (background and foreground progress bar)
            rects = [s for s in slide.shapes if 'AUTO_SHAPE' in str(s.shape_type)]
            if len(rects) < 2:
                print(f"FAIL: Component 4 — Slide {slide_num}: fewer than 2 rectangles")
                continue

            # Background rect (wider or equal) and foreground rect
            # Sort by width descending; first is background, second is foreground
            rects_sorted = sorted(rects, key=lambda s: s.width, reverse=True)
            bg_width = rects_sorted[0].width
            fg_width = rects_sorted[1].width

            if bg_width == 0:
                print(f"FAIL: Component 4 — Slide {slide_num}: background width is 0")
                continue

            actual_pct = fg_width / bg_width
            # For 100% tasks, fg == bg
            if expected_pct == 1.0:
                actual_pct = fg_width / bg_width

            if abs(actual_pct - expected_pct) <= tolerance:
                total_score += per_slide_pts4
                print(f"PASS: Component 4 — Slide {slide_num}: bar {round(actual_pct*100)}% (expected {task['completion_pct']}%)")
            else:
                print(f"FAIL: Component 4 — Slide {slide_num}: bar {round(actual_pct*100)}% (expected {task['completion_pct']}%)")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Final slide is a Blockers slide listing blocked tasks (0.15 pts)
    try:
        last_slide = prs.slides[-1]

        # Check "Blockers" title/heading
        all_text = ''
        for shape in last_slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    all_text += para.text + '\n'

        all_text_lower = all_text.lower()

        if 'blocker' in all_text_lower:
            total_score += 0.05
            print("PASS: Component 5a — 'Blockers' heading found on last slide (0.05 pts)")
        else:
            print("FAIL: Component 5a — 'Blockers' heading not found on last slide")

        # Check blocked tasks are listed
        blocked_found = 0
        for bt in blocked_tasks:
            if bt['task_name'].lower() in all_text_lower:
                blocked_found += 1

        if len(blocked_tasks) > 0 and blocked_found == len(blocked_tasks):
            total_score += 0.10
            print(f"PASS: Component 5b — All {blocked_found}/{len(blocked_tasks)} blocked tasks listed (0.10 pts)")
        elif blocked_found > 0:
            partial = 0.10 * (blocked_found / len(blocked_tasks))
            if partial > 0:
                total_score += partial
            print(f"PARTIAL: Component 5b — {blocked_found}/{len(blocked_tasks)} blocked tasks listed ({round(partial, 3)} pts)")
        else:
            print(f"FAIL: Component 5b — No blocked tasks found on last slide")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {round(total_score, 4)}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
verify_task()
