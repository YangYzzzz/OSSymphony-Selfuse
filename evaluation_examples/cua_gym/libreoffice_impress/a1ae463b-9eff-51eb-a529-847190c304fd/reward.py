"""
Reward Script: Color swatch row on slide 8 of Palette.pptx
Task ID: impress_ndo_056
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Five 3cm x 3cm squares exist on slide 8
  Component 2 (0.35): Squares have correct fill colors in order
  Component 3 (0.20): Text boxes below each square with correct hex codes
  Component 4 (0.15): Text boxes use Liberation Mono 10pt font
"""

import os
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_056'

# Expected colors in order (left to right)
EXPECTED_COLORS = ['264653', '2A9D8F', 'E9C46A', 'F4A261', 'E76F51']
EXPECTED_HEX_TEXTS = ['#264653', '#2A9D8F', '#E9C46A', '#F4A261', '#E76F51']

# 3cm in EMU = 1080000
SQUARE_SIZE_EMU = 1080000
# Tolerance: 5% relative
SIZE_TOL = 0.05

def approx_eq(a, b, tol=0.05):
    """Check approximate equality with relative tolerance."""
    if a == b:
        return True
    if a == 0 or b == 0:
        return abs(a - b) < 50000  # small absolute tolerance for zero
    return abs(a - b) / max(abs(a), abs(b)) <= tol


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 8 slides
    if len(prs.slides) < 8:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 8")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[7]  # 0-indexed: slide 8

    # Collect rectangles (AUTO_SHAPE) and text boxes on slide 8
    # Exclude pre-existing shapes: title placeholder and the "Brand Colors" text box
    rectangles = []
    textboxes = []

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            rectangles.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            # Get text content
            text = ""
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
            # Only consider text boxes that contain hex-like text (task-introduced)
            if text.startswith('#') and len(text) == 7:
                textboxes.append(shape)

    # Sort rectangles by left position (left to right)
    rectangles.sort(key=lambda s: s.left)
    textboxes.sort(key=lambda s: s.left)

    # Component 1: Five 3cm x 3cm squares exist on slide 8 (0.30 points)
    try:
        if len(rectangles) >= 5:
            squares_ok = 0
            for i, rect in enumerate(rectangles[:5]):
                w_ok = approx_eq(rect.width, SQUARE_SIZE_EMU)
                h_ok = approx_eq(rect.height, SQUARE_SIZE_EMU)
                if w_ok and h_ok:
                    squares_ok += 1
                else:
                    print(f"  Square {i}: size ({rect.width}, {rect.height}) != expected ~{SQUARE_SIZE_EMU}")

            if squares_ok == 5:
                print(f"PASS: Component 1 -- All 5 squares are ~3cm x 3cm (0.30 pts)")
                total_score += 0.30
            elif squares_ok >= 3:
                partial = 0.30 * (squares_ok / 5)
                print(f"PARTIAL: Component 1 -- {squares_ok}/5 squares correct size ({partial:.2f} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 1 -- Only {squares_ok}/5 squares have correct size")
        else:
            print(f"FAIL: Component 1 -- Found {len(rectangles)} rectangles, expected 5")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Squares have correct fill colors in order (0.35 points)
    try:
        colors_matched = 0
        if len(rectangles) >= 5:
            for i, rect in enumerate(rectangles[:5]):
                try:
                    fill = rect.fill
                    if fill.type == 1:  # SOLID
                        actual_color = str(fill.fore_color.rgb).upper()
                        expected_color = EXPECTED_COLORS[i].upper()
                        if actual_color == expected_color:
                            colors_matched += 1
                        else:
                            print(f"  Square {i}: color {actual_color} != expected {expected_color}")
                    else:
                        print(f"  Square {i}: fill type is {fill.type}, not SOLID")
                except Exception as e:
                    print(f"  Square {i}: fill error: {e}")

            if colors_matched == 5:
                print(f"PASS: Component 2 -- All 5 colors correct (0.35 pts)")
                total_score += 0.35
            elif colors_matched >= 1:
                partial = 0.35 * (colors_matched / 5)
                print(f"PARTIAL: Component 2 -- {colors_matched}/5 colors correct ({partial:.2f} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 2 -- No colors matched")
        else:
            print(f"FAIL: Component 2 -- Not enough rectangles to check colors")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Text boxes below each square with correct hex codes (0.20 points)
    try:
        texts_matched = 0
        if len(textboxes) >= 5:
            for i, tb in enumerate(textboxes[:5]):
                text = tb.text_frame.text.strip().upper()
                expected = EXPECTED_HEX_TEXTS[i].upper()
                if text == expected:
                    texts_matched += 1
                else:
                    print(f"  TextBox {i}: text '{text}' != expected '{expected}'")

            if texts_matched == 5:
                print(f"PASS: Component 3 -- All 5 hex code labels correct (0.20 pts)")
                total_score += 0.20
            elif texts_matched >= 1:
                partial = 0.20 * (texts_matched / 5)
                print(f"PARTIAL: Component 3 -- {texts_matched}/5 labels correct ({partial:.2f} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 3 -- No hex labels matched")
        else:
            print(f"FAIL: Component 3 -- Found {len(textboxes)} hex text boxes, expected 5")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Text boxes use Liberation Mono 10pt font (0.15 points)
    try:
        font_matched = 0
        if len(textboxes) >= 5:
            for i, tb in enumerate(textboxes[:5]):
                for para in tb.text_frame.paragraphs:
                    for run in para.runs:
                        font_name = run.font.name
                        font_size = run.font.size
                        name_ok = font_name is not None and 'liberation mono' in font_name.lower()
                        # 10pt = 127000 EMU (Pt(10) = 127000)
                        size_ok = font_size is not None and approx_eq(font_size, 127000, tol=0.05)
                        if name_ok and size_ok:
                            font_matched += 1
                        else:
                            print(f"  TextBox {i}: font='{font_name}', size={font_size} (expected Liberation Mono, 127000)")
                        break  # check first run only
                    break  # check first paragraph only

            if font_matched == 5:
                print(f"PASS: Component 4 -- All 5 labels use Liberation Mono 10pt (0.15 pts)")
                total_score += 0.15
            elif font_matched >= 1:
                partial = 0.15 * (font_matched / 5)
                print(f"PARTIAL: Component 4 -- {font_matched}/5 labels have correct font ({partial:.2f} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 4 -- No labels have correct font")
        else:
            print(f"FAIL: Component 4 -- Not enough text boxes to check font")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")

# Main execution
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
