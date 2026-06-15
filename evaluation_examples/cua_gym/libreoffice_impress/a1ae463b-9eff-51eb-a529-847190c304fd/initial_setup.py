"""
Initial Setup: Create a presentation with 8 slides, slide 8 titled 'Brand Colors' with no shapes.
Task ID: impress_ndo_056
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_056'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Palette Design System"
    slide1.placeholders[1].text = "Brand Identity Guidelines v2.3"

    # --- Slide 2: Introduction ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Introduction"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "This design system establishes the visual foundation for our brand across all digital and print touchpoints."
    p2 = body2.add_paragraph()
    p2.text = "Consistency in color usage reinforces brand recognition and trust with our audience."
    p3 = body2.add_paragraph()
    p3.text = "Last updated: March 2026 by the Design Operations team."

    # --- Slide 3: Typography ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Typography Standards"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Primary Typeface: Inter (headings, UI elements)"
    p = body3.add_paragraph()
    p.text = "Secondary Typeface: Merriweather (body text, long-form content)"
    p = body3.add_paragraph()
    p.text = "Monospace: Liberation Mono (code snippets, data tables)"
    p = body3.add_paragraph()
    p.text = "Minimum body size: 14pt for presentations, 11pt for documents"

    # --- Slide 4: Layout Grid ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Layout Grid System"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "12-column grid with 24px gutters for web layouts"
    p = body4.add_paragraph()
    p.text = "8-point spacing system for consistent vertical rhythm"
    p = body4.add_paragraph()
    p.text = "Maximum content width: 1200px for desktop breakpoints"
    p = body4.add_paragraph()
    p.text = "Mobile-first responsive approach with breakpoints at 576px, 768px, 1024px"

    # --- Slide 5: Iconography ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Iconography Guidelines"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Icon style: Outlined, 2px stroke weight, rounded corners"
    p = body5.add_paragraph()
    p.text = "Standard sizes: 16px, 24px, 32px, 48px"
    p = body5.add_paragraph()
    p.text = "All icons must pass WCAG 2.1 contrast requirements against their background"
    p = body5.add_paragraph()
    p.text = "Custom icon library maintained in the shared Figma workspace"

    # --- Slide 6: Photography ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Photography Direction"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Natural lighting preferred; avoid overly processed or filtered images"
    p = body6.add_paragraph()
    p.text = "Authentic representation of diverse communities and environments"
    p = body6.add_paragraph()
    p.text = "Color grading should align with brand palette warm tones"
    p = body6.add_paragraph()
    p.text = "Minimum resolution: 300 DPI for print, 72 DPI for web"

    # --- Slide 7: Accessibility ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Accessibility Standards"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "All color combinations must meet WCAG 2.1 AA contrast ratio (4.5:1 for text)"
    p = body7.add_paragraph()
    p.text = "Never rely on color alone to convey information"
    p = body7.add_paragraph()
    p.text = "Interactive elements require minimum 44x44px touch targets"
    p = body7.add_paragraph()
    p.text = "Alt text required for all decorative and informational images"

    # --- Slide 8: Brand Colors (EMPTY - no shapes, just title) ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add only a title text box at the top
    txBox = slide8.shapes.add_textbox(Cm(2), Cm(0.5), Cm(20), Cm(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Brand Colors"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
