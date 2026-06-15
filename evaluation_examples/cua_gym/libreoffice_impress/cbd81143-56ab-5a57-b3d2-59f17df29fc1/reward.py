"""
Reward Script: Insert comparison table on slide 8 with highlighted better values
Task ID: impress_exec_067
Domain: libreoffice_impress
Scoring:
  Component 1 (0.2): Table exists on slide 8 with 4 rows x 3 columns
  Component 2 (0.2): Header row has correct values (Criteria, Target A, Target B)
  Component 3 (0.3): Data rows have correct values
  Component 4 (0.3): Better values highlighted with #C8E6C9 cell fill
"""

import os
import sys

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_067'


def persist_app_state(domain):
    """Attempt to save any unsaved LibreOffice state."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            import time
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_cell_fill_color(cell):
    """Extract solid fill color from a table cell via XML."""
    try:
        from pptx.oxml.ns import qn
        tc = cell._tc
        tcPr = tc.find(qn('a:tcPr'))
        if tcPr is not None:
            solidFill = tcPr.find(qn('a:solidFill'))
            if solidFill is not None:
                srgbClr = solidFill.find(qn('a:srgbClr'))
                if srgbClr is not None:
                    return srgbClr.get('val').upper()
    except Exception:
        pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 8 slides
    if len(prs.slides) < 8:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 8")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[7]  # 0-indexed, slide 8

    # Find the table on slide 8
    table = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            break

    # Component 1: Table exists with correct dimensions (0.2 points)
    try:
        if table is not None:
            num_rows = len(table.rows)
            num_cols = len(table.columns)
            if num_rows == 4 and num_cols == 3:
                print(f"PASS: Component 1 -- Table found on slide 8 with 4 rows x 3 cols (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 1 -- Table dimensions {num_rows}x{num_cols}, expected 4x3")
        else:
            print(f"FAIL: Component 1 -- No table found on slide 8")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    if table is None:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Header row correct (0.2 points)
    try:
        expected_headers = ['Criteria', 'Target A', 'Target B']
        actual_headers = [table.cell(0, c).text.strip() for c in range(min(len(table.columns), 3))]
        if actual_headers == expected_headers:
            print(f"PASS: Component 2 -- Headers correct: {actual_headers} (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 2 -- Headers {actual_headers}, expected {expected_headers}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Data values correct (0.3 points)
    # 3 data rows x 3 columns = 9 cells, each worth 0.3/9 ~ 0.033
    try:
        expected_data = [
            ['Revenue', '$15M', '$22M'],
            ['Growth Rate', '45%', '28%'],
            ['Asking Price', '$120M', '$180M'],
        ]
        data_correct = 0
        data_total = 9
        for r_idx, expected_row in enumerate(expected_data):
            actual_row_idx = r_idx + 1  # skip header
            if actual_row_idx >= len(table.rows):
                print(f"FAIL: Component 3 -- Missing data row {r_idx + 1}")
                continue
            for c_idx, expected_val in enumerate(expected_row):
                if c_idx >= len(table.columns):
                    continue
                actual_val = table.cell(actual_row_idx, c_idx).text.strip()
                if actual_val == expected_val:
                    data_correct += 1
                else:
                    print(f"FAIL: Component 3 -- Cell({actual_row_idx},{c_idx}) = '{actual_val}', expected '{expected_val}'")

        data_score = 0.3 * (data_correct / data_total)
        if data_correct == data_total:
            print(f"PASS: Component 3 -- All {data_total} data cells correct (0.3 pts)")
            total_score += 0.3
        elif data_correct > 0:
            print(f"PARTIAL: Component 3 -- {data_correct}/{data_total} data cells correct ({data_score:.3f} pts)")
            total_score += data_score
        else:
            print(f"FAIL: Component 3 -- No data cells matched")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Correct cells highlighted with #C8E6C9 fill (0.3 points)
    # Expected highlights:
    #   Revenue: Target B ($22M > $15M) -> cell(1,2)
    #   Growth Rate: Target A (45% > 28%) -> cell(2,1)
    #   Asking Price: Target A ($120M < $180M, lower is better) -> cell(3,1)
    try:
        highlight_color = 'C8E6C9'
        expected_highlighted = [(1, 2), (2, 1), (3, 1)]
        # All other data cells should NOT have the highlight
        expected_not_highlighted = [(1, 1), (2, 2), (3, 2)]

        highlight_correct = 0
        highlight_total = 6  # 3 should be highlighted + 3 should not

        for (r, c) in expected_highlighted:
            if r >= len(table.rows) or c >= len(table.columns):
                continue
            fill = get_cell_fill_color(table.cell(r, c))
            if fill == highlight_color:
                highlight_correct += 1
            else:
                print(f"FAIL: Component 4 -- Cell({r},{c}) '{table.cell(r,c).text}' should have #C8E6C9 fill, found {fill}")

        for (r, c) in expected_not_highlighted:
            if r >= len(table.rows) or c >= len(table.columns):
                continue
            fill = get_cell_fill_color(table.cell(r, c))
            if fill != highlight_color:
                highlight_correct += 1
            else:
                print(f"FAIL: Component 4 -- Cell({r},{c}) '{table.cell(r,c).text}' should NOT have #C8E6C9 fill")

        highlight_score = 0.3 * (highlight_correct / highlight_total)
        if highlight_correct == highlight_total:
            print(f"PASS: Component 4 -- All highlight checks passed (0.3 pts)")
            total_score += 0.3
        elif highlight_correct > 0:
            print(f"PARTIAL: Component 4 -- {highlight_correct}/{highlight_total} highlight checks correct ({highlight_score:.3f} pts)")
            total_score += highlight_score
        else:
            print(f"FAIL: Component 4 -- No highlight checks passed")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
