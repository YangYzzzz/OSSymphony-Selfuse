"""
Initial Setup: M&A Review presentation with 10 slides, slide 8 empty except title
Task ID: impress_exec_067
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_067'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_body(slide, title_text, body_lines):
    """Helper to set title and body content on a slide."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    # Find body placeholder (index 1 typically)
    body_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            body_ph = ph
            break
    if body_ph and body_lines:
        tf = body_ph.text_frame
        tf.clear()
        for i, line in enumerate(body_lines):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                p = tf.add_paragraph()
                p.text = line


def add_title_only(slide, title_text):
    """Set only the title on a slide."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text


def create_initial():
    prs = Presentation()

    # Slide 1: Title Slide (layout 0)
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "M&A Review"
    slide1.placeholders[1].text = "Q4 2025 Strategic Analysis"

    # Slide 2: Executive Summary (layout 1 = Title + Content)
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide2, "Executive Summary", [
        "Three potential acquisition targets identified in the technology sector",
        "Combined deal value estimated at $315M across all targets",
        "Due diligence initiated for two primary candidates",
        "Board approval anticipated by end of Q1 2026",
    ])

    # Slide 3: Market Overview
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide3, "Market Overview", [
        "SaaS market projected to reach $908B by 2030",
        "Enterprise collaboration tools growing at 14.2% CAGR",
        "Key competitors: Salesforce, Microsoft, SAP",
        "Regulatory environment remains favorable for tech M&A",
    ])

    # Slide 4: Target A Profile
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide4, "Target A - CloudSync Analytics", [
        "Founded: 2018 | Headquarters: Austin, TX",
        "Revenue: $15M ARR | Growth Rate: 45% YoY",
        "Employees: 120 | Key Product: Real-time data pipeline platform",
        "Asking Price: $120M (8x revenue multiple)",
        "Key Strength: Proprietary ML-based data integration engine",
    ])

    # Slide 5: Target B Profile
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide5, "Target B - DataVault Solutions", [
        "Founded: 2015 | Headquarters: Seattle, WA",
        "Revenue: $22M ARR | Growth Rate: 28% YoY",
        "Employees: 185 | Key Product: Enterprise data governance suite",
        "Asking Price: $180M (8.2x revenue multiple)",
        "Key Strength: Established Fortune 500 client base",
    ])

    # Slide 6: Financial Analysis
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide6, "Financial Analysis", [
        "Target A NPV: $142M at 12% discount rate",
        "Target B NPV: $198M at 12% discount rate",
        "Synergy potential estimated at $8M-12M annually",
        "Integration costs projected at $15M over 18 months",
        "Payback period: 3.2 years (Target A) / 4.1 years (Target B)",
    ])

    # Slide 7: Risk Assessment
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide7, "Risk Assessment", [
        "Technology integration complexity: Medium (Target A) / Low (Target B)",
        "Key employee retention risk: High for both targets",
        "Regulatory approval timeline: 60-90 days estimated",
        "Customer overlap: 12% with existing portfolio",
        "IP litigation pending: None identified for either target",
    ])

    # Slide 8: Acquisition Comparison - TITLE ONLY, NO CONTENT
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title as a text box since blank layout has no title placeholder
    from pptx.util import Emu
    txBox = slide8.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Acquisition Comparison"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # Slide 9: Integration Timeline
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide9, "Integration Timeline", [
        "Month 1-3: Legal and regulatory approvals",
        "Month 3-6: Technology platform assessment and migration planning",
        "Month 6-9: Core systems integration and data migration",
        "Month 9-12: Full operational integration and optimization",
        "Month 12-18: Post-merger performance review and adjustments",
    ])

    # Slide 10: Recommendations & Next Steps
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide10, "Recommendations & Next Steps", [
        "Proceed with formal LOI for Target A based on superior growth metrics",
        "Initiate parallel due diligence for Target B as backup option",
        "Engage external legal counsel for regulatory filing preparation",
        "Schedule board strategy session for February 15, 2026",
        "Prepare detailed integration playbook by March 1, 2026",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
