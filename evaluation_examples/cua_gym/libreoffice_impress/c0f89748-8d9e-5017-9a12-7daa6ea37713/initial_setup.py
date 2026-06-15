"""
Initial Setup: Fix broken video embed on slide 8
Task ID: impress_fix_028
Domain: libreoffice_impress

Creates a Product_Demo.pptx with 8 slides. Slide 8 has a broken video embed
(black rectangle placeholder). Also creates ~/Desktop/demo_video.mp4 (a real
30-second 1080p video) that the agent should use to re-embed.
"""

import os
import shlex
import subprocess
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_028'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
VIDEO_PATH = f'{WORKDIR}/Desktop/demo_video.mp4'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_demo_video():
    """Create a 30-second 1080p demo video using ffmpeg."""
    os.makedirs(os.path.dirname(VIDEO_PATH), exist_ok=True)
    # Generate a 30-second video with color bars and text overlay
    cmd = (
        'ffmpeg -y -f lavfi -i '
        '"color=c=0x2B5797:s=1920x1080:d=30,drawtext='
        "fontfile=/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf:"
        "text='Product Demo Video':fontcolor=white:fontsize=72:"
        "x=(w-text_w)/2:y=(h-text_h)/2,"
        "drawtext=fontfile=/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf:"
        "text='%{eif\\\\:t\\\\:d} sec':fontcolor=white:fontsize=36:"
        'x=(w-text_w)/2:y=(h+100)/2"'
        ' -f lavfi -i "sine=frequency=440:duration=30"'
        f' -c:v libx264 -preset ultrafast -crf 28 -c:a aac -shortest "{VIDEO_PATH}"'
    )
    result = subprocess.run(cmd, shell=True, capture_output=True, text=True, timeout=120)
    if result.returncode != 0:
        # Fallback: simpler video without text overlays
        cmd_simple = (
            f'ffmpeg -y -f lavfi -i "color=c=0x2B5797:s=1920x1080:d=30" '
            f'-f lavfi -i "sine=frequency=440:duration=30" '
            f'-c:v libx264 -preset ultrafast -crf 28 -c:a aac -shortest "{VIDEO_PATH}"'
        )
        subprocess.run(cmd_simple, shell=True, capture_output=True, text=True, timeout=120)
    print(f'Demo video created: {VIDEO_PATH}')


def create_presentation():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color palette
    DARK_BLUE = RGBColor(0x2B, 0x57, 0x97)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
    ACCENT_GREEN = RGBColor(0x2E, 0x8B, 0x57)

    def add_text_box(slide, left, top, width, height, text, font_size=18,
                     bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT,
                     font_name="Calibri"):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = alignment
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
        return txBox

    # ── Slide 1: Title Slide ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    add_text_box(slide1, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
                 "CloudSync Pro", font_size=48, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, Inches(2), Inches(3.8), Inches(9), Inches(1),
                 "Product Demo Presentation — Q2 2025 Launch", font_size=24,
                 color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, Inches(3), Inches(5.5), Inches(7), Inches(0.6),
                 "Prepared by: Sarah Chen, VP of Product Engineering",
                 font_size=16, color=RGBColor(0xBB, 0xCC, 0xDD),
                 alignment=PP_ALIGN.CENTER)

    # ── Slide 2: Agenda ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide2, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                 "Agenda", font_size=36, bold=True, color=DARK_BLUE)
    agenda_items = [
        "1. Market Overview & Opportunity",
        "2. Product Architecture",
        "3. Key Features & Differentiators",
        "4. Performance Benchmarks",
        "5. Security & Compliance",
        "6. Pricing Tiers",
        "7. Roadmap & Timeline",
        "8. Live Product Demo Video",
    ]
    for i, item in enumerate(agenda_items):
        add_text_box(slide2, Inches(1.2), Inches(1.6 + i * 0.65), Inches(9), Inches(0.6),
                     item, font_size=20, color=DARK_GRAY)

    # ── Slide 3: Market Overview ──
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide3, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                 "Market Overview", font_size=36, bold=True, color=DARK_BLUE)
    market_points = [
        "• Global cloud storage market projected to reach $376B by 2029",
        "• Enterprise file sync and share growing at 24.1% CAGR",
        "• 73% of enterprises plan hybrid cloud adoption by 2026",
        "• Average data breach cost increased to $4.88M in 2024",
        "• Remote workforce driving demand for secure collaboration tools",
    ]
    for i, point in enumerate(market_points):
        add_text_box(slide3, Inches(1.0), Inches(1.8 + i * 0.85), Inches(10), Inches(0.7),
                     point, font_size=18, color=DARK_GRAY)

    # ── Slide 4: Product Architecture ──
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide4, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                 "Product Architecture", font_size=36, bold=True, color=DARK_BLUE)
    # Simplified architecture boxes
    components = [
        ("Client SDK", 1.0, 2.0, 3.0, 1.5),
        ("API Gateway", 5.0, 2.0, 3.0, 1.5),
        ("Storage Engine", 9.0, 2.0, 3.0, 1.5),
        ("Auth Service", 1.0, 4.5, 3.0, 1.5),
        ("Sync Orchestrator", 5.0, 4.5, 3.0, 1.5),
        ("Analytics Pipeline", 9.0, 4.5, 3.0, 1.5),
    ]
    from pptx.enum.shapes import MSO_SHAPE
    for name, l, t, w, h in components:
        shape = slide4.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_BLUE
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(16)
        run.font.color.rgb = WHITE
        run.font.bold = True

    # ── Slide 5: Key Features ──
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide5, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                 "Key Features & Differentiators", font_size=36, bold=True, color=DARK_BLUE)
    features = [
        ("Real-Time Collaboration", "Multi-user editing with conflict-free replicated data types (CRDTs)"),
        ("Zero-Knowledge Encryption", "End-to-end encryption where only users hold decryption keys"),
        ("Smart Sync Engine", "Delta sync with intelligent caching reduces bandwidth by up to 85%"),
        ("Cross-Platform Support", "Native apps for Windows, macOS, Linux, iOS, and Android"),
    ]
    for i, (title, desc) in enumerate(features):
        add_text_box(slide5, Inches(1.0), Inches(1.6 + i * 1.3), Inches(10), Inches(0.5),
                     title, font_size=22, bold=True, color=ACCENT_GREEN)
        add_text_box(slide5, Inches(1.3), Inches(2.1 + i * 1.3), Inches(10), Inches(0.5),
                     desc, font_size=16, color=DARK_GRAY)

    # ── Slide 6: Performance Benchmarks ──
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide6, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                 "Performance Benchmarks", font_size=36, bold=True, color=DARK_BLUE)
    # Add a table
    rows, cols = 5, 4
    tbl_shape = slide6.shapes.add_table(rows, cols, Inches(1.0), Inches(1.8),
                                         Inches(11), Inches(3.5))
    table = tbl_shape.table
    headers = ["Metric", "CloudSync Pro", "Competitor A", "Competitor B"]
    data = [
        ["Upload Speed (1GB)", "45 sec", "72 sec", "68 sec"],
        ["Sync Latency", "< 200ms", "1.2 sec", "800ms"],
        ["Concurrent Users", "50,000+", "10,000", "25,000"],
        ["Storage Efficiency", "94.2%", "78.5%", "82.1%"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(13)

    # ── Slide 7: Pricing ──
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide7, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                 "Pricing Tiers", font_size=36, bold=True, color=DARK_BLUE)
    tiers = [
        ("Starter", "$9.99/mo", "50 GB Storage\n5 Users\nBasic Support"),
        ("Professional", "$24.99/mo", "500 GB Storage\n25 Users\nPriority Support"),
        ("Enterprise", "$79.99/mo", "Unlimited Storage\nUnlimited Users\n24/7 Support"),
    ]
    for i, (name, price, desc) in enumerate(tiers):
        left = Inches(1.0 + i * 4.0)
        shape = slide7.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(3.5), Inches(4.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_GRAY
        shape.line.color.rgb = DARK_BLUE
        add_text_box(slide7, left + Inches(0.3), Inches(2.0), Inches(3), Inches(0.7),
                     name, font_size=24, bold=True, color=DARK_BLUE,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide7, left + Inches(0.3), Inches(2.8), Inches(3), Inches(0.7),
                     price, font_size=32, bold=True, color=ACCENT_GREEN,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide7, left + Inches(0.3), Inches(3.8), Inches(3), Inches(2),
                     desc, font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # ── Slide 8: Video Demo (BROKEN) ──
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide8, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                 "Live Product Demo", font_size=36, bold=True, color=DARK_BLUE)
    add_text_box(slide8, Inches(1.0), Inches(1.3), Inches(10), Inches(0.6),
                 "Watch CloudSync Pro in action — file sync, real-time collaboration, and admin dashboard",
                 font_size=16, color=DARK_GRAY)

    # Create a black rectangle to simulate broken video embed
    black_rect = slide8.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2.0), Inches(2.2), Inches(9.333), Inches(5.0)
    )
    black_rect.fill.solid()
    black_rect.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    black_rect.line.fill.background()

    # Add a play button triangle on top of the black rectangle
    play_btn = slide8.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(6.0), Inches(4.0), Inches(1.2), Inches(1.2)
    )
    play_btn.fill.solid()
    play_btn.fill.fore_color.rgb = RGBColor(0x80, 0x80, 0x80)
    play_btn.line.fill.background()
    # Rotate to point right (play icon)
    play_btn.rotation = 90.0

    # Add a note indicating the video is broken
    slide8.notes_slide.notes_text_frame.text = (
        "NOTE: The embedded video link is broken. The original video file "
        "is located at ~/Desktop/demo_video.mp4. Please re-embed and set to autoplay."
    )

    prs.save(OUTPUT)
    print(f'Initial presentation created: {OUTPUT}')


create_demo_video()
create_presentation()

# GUI-ready startup
launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')
