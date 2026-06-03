"""
Reward Script: Re-embed video on slide 8 and set autoplay
Task ID: impress_fix_028
Domain: libreoffice_impress
Scoring:
  Component 1 (0.35): Slide 8 contains a MEDIA shape (properly embedded video)
  Component 2 (0.25): Video file exists in the PPTX ZIP archive (ppt/media/)
  Component 3 (0.25): Auto-play timing configured with playFrom command
  Component 4 (0.15): Broken placeholder shapes (Rectangle + Triangle) removed from slide 8
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_028'


def persist_app_state(domain):
    """Send Ctrl+S to save any unsaved GUI edits."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            import time
            time.sleep(1.0)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Check the file can be loaded by python-pptx
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Must have at least 8 slides
    if len(prs.slides) < 8:
        print(f"CRITICAL: Expected at least 8 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide8 = prs.slides[7]  # 0-indexed

    # Component 1: Slide 8 contains a MEDIA shape (0.35 points)
    # In initial_env, slide 8 has no MEDIA shape (only TextBox + AutoShape placeholders).
    # In golden_env, the broken placeholders are replaced with a real embedded video (MEDIA type 16).
    try:
        media_shapes = [s for s in slide8.shapes if s.shape_type == MSO_SHAPE_TYPE.MEDIA]
        if len(media_shapes) > 0:
            media_shape = media_shapes[0]
            print(f"PASS: Component 1 -- Slide 8 has MEDIA shape: '{media_shape.name}' (0.35 pts)")
            total_score += 0.35
        else:
            shape_types = [(s.name, str(s.shape_type)) for s in slide8.shapes]
            print(f"FAIL: Component 1 -- No MEDIA shape on slide 8. Shapes: {shape_types}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Video file exists in the PPTX ZIP archive (0.25 points)
    # In initial_env, there are no media/video files in the ZIP.
    # In golden_env, ppt/media/video1.mp4 (or similar) should exist.
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            video_files = [n for n in zf.namelist() if n.startswith('ppt/media/') and n.endswith('.mp4')]
            if len(video_files) > 0:
                # Verify the video has non-trivial size (at least 1KB)
                video_info = zf.getinfo(video_files[0])
                if video_info.file_size > 1024:
                    print(f"PASS: Component 2 -- Video file found in ZIP: {video_files[0]} "
                          f"(size: {video_info.file_size} bytes) (0.25 pts)")
                    total_score += 0.25
                else:
                    print(f"FAIL: Component 2 -- Video file {video_files[0]} is too small "
                          f"({video_info.file_size} bytes), likely not a real video")
            else:
                all_media = [n for n in zf.namelist() if n.startswith('ppt/media/')]
                print(f"FAIL: Component 2 -- No .mp4 video file in ZIP. Media files: {all_media}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Auto-play timing configured (0.25 points)
    # In initial_env, there is no p:timing element on slide 8.
    # In golden_env, slide 8 has p:timing with a playFrom(0) command for automatic playback.
    try:
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide8.xml') as f:
                slide_xml = f.read().decode('utf-8')

            root = ET.fromstring(slide_xml)
            timing = root.find('.//p:timing', ns)
            if timing is not None:
                timing_xml = ET.tostring(timing, encoding='unicode')
                # Check for playFrom command which indicates auto-play
                if 'playFrom' in timing_xml:
                    print(f"PASS: Component 3 -- Auto-play timing with playFrom command found (0.25 pts)")
                    total_score += 0.25
                else:
                    print(f"FAIL: Component 3 -- p:timing element exists but no playFrom command found")
            else:
                print(f"FAIL: Component 3 -- No p:timing element found on slide 8 (no autoplay)")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Broken placeholder shapes removed (0.15 points)
    # In initial_env, slide 8 has 4 shapes including a Rectangle and Isosceles Triangle
    # that simulate a broken video (black rectangle + play button triangle).
    # In golden_env, those placeholder shapes should be gone (replaced by the real video).
    try:
        auto_shapes = []
        for s in slide8.shapes:
            if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                auto_shapes.append(s.name)

        # Check that neither the Rectangle nor Triangle placeholder exists
        has_rectangle = any('Rectangle' in name for name in auto_shapes)
        has_triangle = any('Triangle' in name for name in auto_shapes)

        if not has_rectangle and not has_triangle:
            print(f"PASS: Component 4 -- Broken placeholder shapes removed from slide 8 (0.15 pts)")
            total_score += 0.15
        else:
            remaining = [n for n in auto_shapes if 'Rectangle' in n or 'Triangle' in n]
            print(f"FAIL: Component 4 -- Placeholder shapes still present: {remaining}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
