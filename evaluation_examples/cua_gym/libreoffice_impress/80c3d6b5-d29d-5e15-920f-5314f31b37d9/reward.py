"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 77 has a bright photo tagged as “Picture 1,” but the rest of my deck uses a monochrome look. In LibreOffice Impress, how do I convert that specific image to full grayscale—i.e., drop its color saturation to 0%—while leaving everything else untouched?
Generated: 2025-09-10 16:19:12
Status: success
Model: azure-o3
Total Steps: 7
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Namespace used by PresentationML for drawing elements
AE_NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

def detect_grayscale(blip_element):
    """Return True if the <a:blip> element carries a grayscale or 0-% saturation effect."""
    # 1) Explicit <a:grayscl/> tag → definitely grayscale
    if blip_element.find(AE_NS + 'grayscl') is not None:
        return True

    # 2) A <a:satMod val="0"/> anywhere under <a:blip> also indicates 0-% saturation
    for elem in blip_element.iter():
        if elem.tag == AE_NS + 'satMod':
            val = elem.get('val')
            try:
                if int(val) == 0:
                    return True
            except (TypeError, ValueError):
                pass  # Non-numeric or missing value ⇒ ignore
    return False

def verify_task(file_path):
    """Verify that on slide 77 the image named ‘Picture 1’ was converted to grayscale, while the rest remain unchanged."""
    print(f"Verifying presentation: {file_path}")

    total_score = 0.0
    max_score   = 1.0  # upper bound

    # ------------------------------------------------------------------
    # Prerequisite: file must exist & be loadable (no points awarded)
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File not found – task failed")
        return 0.0
    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded successfully with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Could not load presentation: {e}")
        return 0.0

    # ------------------------------------------------------------------
    # Requirement 1 – locate slide 77 and find the target image
    # ------------------------------------------------------------------
    if len(prs.slides) < 77:
        print("✗ Slide 77 does not exist – task failed")
        return 0.0

    slide77 = prs.slides[76]  # zero-based index

    target_shape          = None  # the ‘Picture 1’ shape
    target_is_grayscale   = False
    other_gray_on_slide77 = False  # any other picture on slide 77 turned gray?

    for shape in slide77.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            blip = shape._element.find('.//' + AE_NS + 'blip')
            if blip is None:
                continue  # shouldn’t happen, but be safe
            is_gray = detect_grayscale(blip)

            if shape.name.strip().lower() == 'picture 1':
                target_shape        = shape
                target_is_grayscale = is_gray
            else:
                if is_gray:
                    other_gray_on_slide77 = True  # shouldn’t be gray per instructions

    if target_shape is None:
        print("✗ No image named 'Picture 1' found on slide 77")
        return 0.0

    print("✓ Located 'Picture 1' on slide 77 (0.3 points)")
    total_score += 0.3  # only if picture is present

    # ------------------------------------------------------------------
    # Requirement 2 – ‘Picture 1’ must be grayscale
    # ------------------------------------------------------------------
    if target_is_grayscale:
        print("✓ 'Picture 1' is fully grayscale (0.7 points)")
        total_score += 0.7
    else:
        print("✗ 'Picture 1' is NOT grayscale – main requirement failed")

    # (Optional) small penalty if other images on slide 77 were incorrectly turned gray
    if other_gray_on_slide77 and total_score > 0.7:
        print("ℹ︎ Other pictures on slide 77 are grayscale – subtracting 0.2 points")
        total_score -= 0.2

    # ------------------------------------------------------------------
    final_score = max(0.0, min(total_score, max_score))
    print(f"REWARD: {final_score}")
    return final_score

# ----------------------------------------------------------------------
# Execute verification when run as a script
# ----------------------------------------------------------------------
if __name__ == '__main__':
    verify_task('/home/user/slide_77_has_a_bright_photo_tagged_as_picture_1_but_the_rest_of_my_deck_uses_a_monochrome_look_in_li_golden.pptx')
