"""
Initial Setup: Create a 6-slide presentation with Wipe transition (From Left) on slide 3.
Task ID: impress_tm_017
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import zipfile
import shutil
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_017'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

# Namespaces for OOXML
NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q2 2025 Marketing Strategy"
    slide1.placeholders[1].text = "Prepared by the Digital Marketing Team"

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Campaign Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Social media engagement grew 34% in Q1"
    p2 = body2.add_paragraph()
    p2.text = "Email open rates exceeded industry average by 12%"
    p2.level = 0
    p3 = body2.add_paragraph()
    p3.text = "New customer acquisition cost reduced to $18.50"
    p3.level = 0
    p4 = body2.add_paragraph()
    p4.text = "Brand awareness score improved from 62 to 78"
    p4.level = 0

    # --- Slide 3: Key Metrics (this slide gets the Wipe transition) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Key Performance Metrics"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Website Traffic: 2.4M monthly visitors (+18%)"
    for item in [
        "Conversion Rate: 3.8% (target: 4.0%)",
        "Customer Lifetime Value: $1,245",
        "Net Promoter Score: 72",
        "Social Media Followers: 458K across platforms",
        "Content Engagement: 12.5% average interaction rate",
    ]:
        para = body3.add_paragraph()
        para.text = item
        para.level = 0

    # --- Slide 4: Regional Performance ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Regional Performance Breakdown"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "North America: $4.2M revenue (42% of total)"
    for item in [
        "Europe: $2.8M revenue (28% of total)",
        "Asia-Pacific: $1.9M revenue (19% of total)",
        "Latin America: $0.7M revenue (7% of total)",
        "Middle East & Africa: $0.4M revenue (4% of total)",
    ]:
        para = body4.add_paragraph()
        para.text = item
        para.level = 0

    # --- Slide 5: Action Items ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Q2 Action Items"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Launch TikTok campaign by April 15"
    for item in [
        "Hire 2 additional content creators",
        "Redesign landing page for mobile optimization",
        "Partner with 5 micro-influencers in lifestyle niche",
        "Implement A/B testing framework for email campaigns",
    ]:
        para = body5.add_paragraph()
        para.text = item
        para.level = 0

    # --- Slide 6: Thank You ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[0])
    slide6.shapes.title.text = "Thank You"
    slide6.placeholders[1].text = "Questions? Contact: marketing@globalreach.com"

    # Save the base presentation first
    prs.save(OUTPUT)
    print(f"Base presentation saved: {OUTPUT}")

    # Now add the Wipe transition (From Left) to slide 3 via XML manipulation
    add_wipe_transition(OUTPUT, slide_number=3, direction='l')
    print(f"Initial file created with Wipe (From Left) on slide 3: {OUTPUT}")

    # Launch LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print("GUI_READY: launched LibreOffice Impress with DISPLAY=:0")


def add_wipe_transition(pptx_path, slide_number, direction='l'):
    """
    Add a Wipe transition to a specific slide.
    direction: 'l' = from left (default), 'r' = from right, 'u' = from top, 'd' = from bottom
    """
    tmp_path = pptx_path + '.tmp'
    shutil.copy(pptx_path, tmp_path)

    for prefix, uri in NSMAP.items():
        ET.register_namespace(prefix, uri)

    with zipfile.ZipFile(tmp_path, 'r') as zin:
        with zipfile.ZipFile(pptx_path, 'w') as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == f'ppt/slides/slide{slide_number}.xml':
                    root = ET.fromstring(data)
                    # Remove any existing transition
                    for tr in root.findall('{http://schemas.openxmlformats.org/presentationml/2006/main}transition'):
                        root.remove(tr)

                    # Create transition element: <p:transition><p:wipe dir="l"/></p:transition>
                    p_ns = NSMAP['p']
                    transition = ET.SubElement(root, f'{{{p_ns}}}transition')
                    wipe = ET.SubElement(transition, f'{{{p_ns}}}wipe')
                    if direction != 'l':  # 'l' is default, but let's be explicit
                        wipe.set('dir', direction)
                    else:
                        wipe.set('dir', direction)

                    data = ET.tostring(root, xml_declaration=True, encoding='UTF-8')
                zout.writestr(item, data)

    os.remove(tmp_path)


create_initial()
