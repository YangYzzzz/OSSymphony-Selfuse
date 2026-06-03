"""
FINAL REWARD SCRIPT - SUCCESS
Task: Insert a cross-reference to Heading 2 'Results' that shows the reference text.
Generated: 2025-10-17 18:01:54
Status: success
Model: azure-o3
Total Steps: 19
"""

import os
from pptx import Presentation


def verify_task(file_path: str) -> float:
    """Verify that a presentation contains
    1) a dedicated heading-style slide whose title text is exactly
       the word 'Results' (Heading 2 in the original writer task)
    2) at least one additional occurrence of the *same* text on a
       different slide (acting as the cross-reference that shows the
       reference text).

    Scoring (progressive – max 1.0):
        0.5  – heading slide detected
        0.5  – separate slide also contains the text (cross-reference)
    """
    print(f"Loading presentation: {file_path}")

    # Preliminary checks – file must exist and be a .pptx
    if not os.path.exists(file_path):
        print("✗ File does not exist")
        return 0.0
    if not file_path.lower().endswith(".pptx"):
        print("✗ File is not a PPTX presentation")
        return 0.0

    # Try to open the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Unable to load PPTX: {e}")
        return 0.0

    # --- gather all occurrences of the exact text 'Results' (case-insensitive) ---
    occurrences = []  # tuples -> (slide_index, shape_index)

    for s_idx, slide in enumerate(prs.slides):
        for sh_idx, shape in enumerate(slide.shapes):
            # Only inspect shapes that actually contain text
            if not getattr(shape, "has_text_frame", False):
                continue

            # Extract text exactly as displayed (including line-breaks)
            full_text = shape.text or ""
            if not full_text.strip():
                continue

            # Split on newlines so that bullet items and title lines are
            # considered separately.
            for line in full_text.split("\n"):
                if line.strip().lower() == "results":
                    occurrences.append((s_idx, sh_idx))

    print(f"Found {len(occurrences)} occurrence(s) of the word 'Results'.")

    # No occurrences – task definitely not completed
    if not occurrences:
        print("✗ Text 'Results' not found anywhere – task failed")
        return 0.0

    # ------------------------------------------------------------------
    # Requirement 1 : Heading-style slide containing *only* the word
    #                 'Results' – interpreted here as either a title
    #                 placeholder or a slide that contains very little
    #                 other content (<= 4 shapes).
    # ------------------------------------------------------------------
    heading_slide_found = False
    for slide_idx, _ in occurrences:
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            if shape.text.strip().lower() != "results":
                continue
            # Heuristic: treat as heading if either (a) it's a title
            # placeholder or (b) the slide is minimalistic.
            is_title_placeholder = False
            try:
                if shape.is_placeholder and shape.placeholder_format.type == 1:  # TITLE placeholder
                    is_title_placeholder = True
            except Exception:
                pass

            if is_title_placeholder or len(slide.shapes) <= 4:
                heading_slide_found = True
                break
        if heading_slide_found:
            break

    score = 0.0
    if heading_slide_found:
        print("✓ Heading slide titled 'Results' detected (0.5 points)")
        score += 0.5
    else:
        print("✗ No dedicated heading slide titled 'Results' found (0 points)")

    # ------------------------------------------------------------------
    # Requirement 2 : At least one *other* slide must contain the word
    #                 'Results' – representing the cross-reference.
    # ------------------------------------------------------------------
    slides_with_results = {slide_idx for slide_idx, _ in occurrences}
    if len(slides_with_results) >= 2:
        print("✓ 'Results' also appears on another slide (cross-reference) (0.5 points)")
        score += 0.5
    else:
        print("✗ 'Results' does not appear on a second slide – cross-reference missing (0 points)")

    # Cap to 1.0 just in case
    final_score = min(score, 1.0)
    print(f"Total score: {final_score}")
    return final_score


if __name__ == "__main__":
    FILE_PATH = "/home/user/insert_a_cross_reference_to_heading_2_results_that_shows_the_reference_text.pptx"
    reward = verify_task(FILE_PATH)
    print(f"REWARD: {reward}")

