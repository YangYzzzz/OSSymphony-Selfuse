"""
Initial Setup: Build a 9-slide sales deal pitch presentation
Task ID: impress_sales_093
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_093'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def add_blank_slide_with_title(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(
        prs,
        "DealSummary Pitch Deck",
        "Prepared by Nexus Solutions Inc.\nQ4 2025 Enterprise Offering"
    )

    # Slide 2: Executive Overview
    add_content_slide(prs, "Executive Overview", [
        "Nexus Solutions provides end-to-end cloud infrastructure for mid-market enterprises.",
        "Our platform reduces IT operational costs by an average of 34% within 12 months.",
        "Over 2,400 enterprise clients across North America and EMEA.",
        "Named a Gartner Magic Quadrant Leader for three consecutive years.",
    ])

    # Slide 3: Client Requirements
    add_content_slide(prs, "Client Requirements", [
        "Migrate 85 on-premise servers to managed cloud within 6 months",
        "Achieve 99.95% uptime SLA across all production workloads",
        "Implement role-based access controls for 1,200 employees",
        "Full disaster recovery with <15 minute RTO for critical systems",
        "Compliance with SOC 2 Type II and HIPAA requirements",
    ])

    # Slide 4: Proposed Solution Architecture
    add_content_slide(prs, "Proposed Solution Architecture", [
        "Multi-region deployment across US-East, US-West, and EU-Central",
        "Kubernetes-based container orchestration with auto-scaling",
        "Integrated CI/CD pipeline with automated security scanning",
        "Centralized logging and monitoring via Nexus ObserveHub",
        "Encrypted data at rest (AES-256) and in transit (TLS 1.3)",
    ])

    # Slide 5: Implementation Timeline
    add_content_slide(prs, "Implementation Timeline", [
        "Phase 1 (Weeks 1-4): Discovery, architecture review, environment setup",
        "Phase 2 (Weeks 5-12): Migration of non-critical workloads",
        "Phase 3 (Weeks 13-20): Production migration with zero-downtime cutover",
        "Phase 4 (Weeks 21-24): Optimization, training, and handoff",
        "Ongoing: 24/7 managed support with dedicated account team",
    ])

    # Slide 6: Pricing Overview
    slide6 = add_blank_slide_with_title(prs, "Pricing Overview")
    # Add a simple table for pricing tiers
    tbl_shape = slide6.shapes.add_table(
        4, 3, Inches(1), Inches(1.5), Inches(10), Inches(3)
    )
    tbl = tbl_shape.table
    headers = ["Tier", "Monthly Cost", "Included Resources"]
    tier_data = [
        ["Standard", "$18,500/mo", "50 vCPUs, 200 GB RAM, 5 TB storage"],
        ["Professional", "$34,750/mo", "120 vCPUs, 500 GB RAM, 15 TB storage"],
        ["Enterprise", "$52,000/mo", "250 vCPUs, 1 TB RAM, 40 TB storage + DR"],
    ]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row in enumerate(tier_data, 1):
        for c, val in enumerate(row):
            tbl.cell(r, c).text = val

    # Slide 7: Customer Success Stories
    add_content_slide(prs, "Customer Success Stories", [
        "Meridian Health Systems: 42% cost reduction, zero downtime in 18 months",
        "Atlas Financial Group: Migrated 200+ servers in under 90 days",
        "Pinnacle Retail Corp: Achieved PCI DSS compliance within first quarter",
        "Horizon Education Network: Scaled to support 50,000 concurrent users",
    ])

    # Slide 8: Support & SLA Guarantees
    add_content_slide(prs, "Support & SLA Guarantees", [
        "Dedicated Technical Account Manager assigned to every Enterprise client",
        "99.95% uptime SLA with financial credits for any violation",
        "15-minute response time for Severity 1 incidents (24/7/365)",
        "Quarterly business reviews with executive leadership",
        "Access to Nexus Academy training platform for client teams",
    ])

    # Slide 9: Next Steps
    add_content_slide(prs, "Next Steps", [
        "1. Schedule technical deep-dive with client engineering team",
        "2. Finalize scope of work and migration priority list",
        "3. Execute mutual NDA and Master Services Agreement",
        "4. Kick off Phase 1 discovery within 10 business days of signing",
        "Contact: Sarah Mitchell, VP Enterprise Sales — sarah.mitchell@nexussolutions.com",
    ])

    # NOTE: No slide 10 — the agent task is to create it
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
