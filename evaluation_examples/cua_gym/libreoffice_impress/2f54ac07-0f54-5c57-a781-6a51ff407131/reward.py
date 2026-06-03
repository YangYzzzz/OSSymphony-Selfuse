"""
Reward Script: Build a complete deal summary slide on slide 10
Task ID: impress_sales_093
Domain: libreoffice_impress
Scoring:
  C1 (0.15): Slide 10 exists (pres has >=10 slides)
  C2 (0.25): Table on slide 10 with correct structure (5 rows x 4 cols, headers)
  C3 (0.15): Table has 3 populated line-item rows
  C4 (0.15): Total row has #2B6CB0 background fill and bold text
  C5 (0.10): Payment terms text present on slide 10
  C6 (0.05): 'Authorized By:' label present on slide 10
  C7 (0.05): Signature line shape (AutoShape/rectangle) on slide 10
  C8 (0.10): Date text present on slide 10
"""

import os
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_093'


def persist_app_state(domain: str):
    """Attempt to save any unsaved GUI state."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide 10 exists (0.15 points)
    # Initial has 9 slides, golden has 10. This is a task-introduced change.
    try:
        num_slides = len(prs.slides)
        if num_slides >= 10:
            print(f"PASS: Component 1 — Slide 10 exists (slide count: {num_slides}) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — Expected >=10 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Gate: If no slide 10, remaining checks are meaningless
    if len(prs.slides) < 10:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    slide10 = prs.slides[9]

    # Component 2: Table on slide 10 with correct structure (0.25 points)
    # 5 rows x 4 cols, with headers: Item, Description, Qty, Total
    try:
        table_shape = None
        for s in slide10.shapes:
            if s.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_shape = s
                break

        if table_shape is None:
            print("FAIL: Component 2 — No table found on slide 10")
        else:
            table = table_shape.table
            nrows = len(table.rows)
            ncols = len(table.columns)

            # Check dimensions: need at least 5 rows and 4 columns
            dim_ok = nrows >= 5 and ncols >= 4

            # Check header row
            expected_headers = ['item', 'description', 'qty', 'total']
            actual_headers = [table.cell(0, c).text.strip().lower() for c in range(min(ncols, 4))]
            headers_ok = actual_headers == expected_headers

            if dim_ok and headers_ok:
                print(f"PASS: Component 2 — Table {nrows}x{ncols} with correct headers {actual_headers} (0.25 pts)")
                total_score += 0.25
            elif dim_ok:
                print(f"PARTIAL: Component 2 — Table dimensions ok ({nrows}x{ncols}) but headers wrong: {actual_headers}")
                total_score += 0.10
            else:
                print(f"FAIL: Component 2 — Table dimensions {nrows}x{ncols}, expected >=5x4. Headers: {actual_headers}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Table has 3 populated line-item rows (0.15 points)
    # Rows 1-3 should have non-empty content in Item and Total columns
    try:
        if table_shape is not None:
            table = table_shape.table
            populated_items = 0
            for r in range(1, min(len(table.rows), 4)):  # rows 1,2,3
                item_text = table.cell(r, 0).text.strip()
                total_text = table.cell(r, min(len(table.columns)-1, 3)).text.strip()
                if item_text and total_text:
                    populated_items += 1

            if populated_items >= 3:
                print(f"PASS: Component 3 — {populated_items} populated line-item rows (0.15 pts)")
                total_score += 0.15
            elif populated_items >= 1:
                partial = round(0.15 * populated_items / 3, 2)
                print(f"PARTIAL: Component 3 — Only {populated_items}/3 line items populated ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 3 — No populated line-item rows found")
        else:
            print("FAIL: Component 3 — No table to check")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Total row has #2B6CB0 background and bold text (0.15 points)
    # Check the last row of the table for blue background fill and bold text
    try:
        if table_shape is not None:
            table = table_shape.table
            last_row = len(table.rows) - 1

            # Check background fill via XML
            ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            blue_fill_count = 0
            bold_found = False

            for c in range(min(len(table.columns), 4)):
                cell = table.cell(last_row, c)
                tc_el = cell._tc
                # Look for solidFill with srgbClr
                fill_el = tc_el.find(f'.//{{{ns_a}}}solidFill/{{{ns_a}}}srgbClr')
                if fill_el is not None:
                    val = fill_el.get('val', '').upper()
                    if val == '2B6CB0':
                        blue_fill_count += 1

                # Check for bold text in this cell
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.bold is True:
                            bold_found = True

            # Need at least some cells with blue fill AND bold text somewhere in total row
            has_blue = blue_fill_count >= 2  # At least 2 cells with blue bg
            has_bold = bold_found

            if has_blue and has_bold:
                print(f"PASS: Component 4 — Total row: {blue_fill_count} cells with #2B6CB0 bg, bold text found (0.15 pts)")
                total_score += 0.15
            elif has_blue:
                print(f"PARTIAL: Component 4 — Blue bg found ({blue_fill_count} cells) but no bold text (0.08 pts)")
                total_score += 0.08
            elif has_bold:
                print(f"PARTIAL: Component 4 — Bold text found but no #2B6CB0 bg (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 4 — No blue bg (found {blue_fill_count}) and no bold text")
        else:
            print("FAIL: Component 4 — No table to check")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Payment terms text on slide 10 (0.10 points)
    # Expected: "Net 30, 10% discount for annual commitment" or similar
    try:
        payment_found = False
        for s in slide10.shapes:
            if hasattr(s, 'text') and s.text:
                text_lower = s.text.lower()
                if 'net 30' in text_lower and 'discount' in text_lower:
                    payment_found = True
                    print(f"PASS: Component 5 — Payment terms found: '{s.text[:60]}' (0.10 pts)")
                    total_score += 0.10
                    break
        if not payment_found:
            print("FAIL: Component 5 — Payment terms text not found on slide 10")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: 'Authorized By:' label on slide 10 (0.05 points)
    try:
        auth_found = False
        for s in slide10.shapes:
            if hasattr(s, 'text') and s.text:
                if 'authorized by' in s.text.lower():
                    auth_found = True
                    print(f"PASS: Component 6 — 'Authorized By:' label found (0.05 pts)")
                    total_score += 0.05
                    break
        if not auth_found:
            print("FAIL: Component 6 — 'Authorized By:' label not found on slide 10")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Signature line shape on slide 10 (0.05 points)
    # Looking for an AutoShape (rectangle/line) that serves as a signature line
    try:
        sig_line_found = False
        for s in slide10.shapes:
            # Look for auto shapes (rectangles, lines) that are thin (signature line)
            if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or s.shape_type == MSO_SHAPE_TYPE.FREEFORM:
                # A signature line is typically a thin rectangle or line
                if s.height < s.width and s.height < 100000:  # thin shape, less than ~1 inch tall
                    sig_line_found = True
                    print(f"PASS: Component 7 — Signature line shape found (name={s.name}, h={s.height}) (0.05 pts)")
                    total_score += 0.05
                    break
        if not sig_line_found:
            # Also check for LINE shapes
            for s in slide10.shapes:
                if 'LINE' in str(s.shape_type) or (s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE):
                    sig_line_found = True
                    print(f"PASS: Component 7 — AutoShape found as signature line (name={s.name}) (0.05 pts)")
                    total_score += 0.05
                    break
        if not sig_line_found:
            print("FAIL: Component 7 — No signature line shape found on slide 10")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    # Component 8: Date text on slide 10 (0.10 points)
    # Looking for a date string (any reasonable date format)
    try:
        import re
        date_found = False
        date_patterns = [
            r'\b\d{1,2}/\d{1,2}/\d{2,4}\b',           # MM/DD/YYYY
            r'\b\d{4}-\d{2}-\d{2}\b',                   # YYYY-MM-DD
            r'\b(?:jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)\w*\s+\d{1,2},?\s+\d{4}\b',  # Month DD, YYYY
            r'\b\d{1,2}\s+(?:jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)\w*\s+\d{4}\b',     # DD Month YYYY
        ]
        for s in slide10.shapes:
            if hasattr(s, 'text') and s.text:
                for pattern in date_patterns:
                    if re.search(pattern, s.text.strip(), re.IGNORECASE):
                        date_found = True
                        print(f"PASS: Component 8 — Date found: '{s.text.strip()[:40]}' (0.10 pts)")
                        total_score += 0.10
                        break
                if date_found:
                    break
        if not date_found:
            print("FAIL: Component 8 — No date text found on slide 10")
    except Exception as e:
        print(f"ERROR: Component 8 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
