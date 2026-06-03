"""
FINAL REWARD SCRIPT - SUCCESS
Task: Apply no borders to Table 1 and set cell padding to 0.2 cm.
Generated: 2025-10-17 10:10:05
Status: success
Model: azure-o3
Total Steps: 4
"""

from pptx import Presentation
import os
from lxml import etree

# -----------------------------------------------------------------------------
# Reward Script : Verify "Apply no borders to Table 1 and set cell padding to 0.2 cm"
# -----------------------------------------------------------------------------
# Scoring Rules
#   0.5 pts – Table-1 has ALL borders removed (width == 0)
#   0.5 pts – ALL cell paddings ≈ 0.2 cm (±0.1 cm tolerance)
#   1.0 pts – Both conditions satisfied
# -----------------------------------------------------------------------------
# NOTE: 1 cm = 360 000 EMU in the PPTX specification
# -----------------------------------------------------------------------------

def verify_task(file_path: str) -> float:
    """Return a progressive score (0.0-1.0) based on task completion."""

    print(f"Analyzing presentation: {file_path}")
    total_score = 0.0
    max_score = 1.0

    # ---------------------------------------------------------------------
    # 0. Prerequisite: file must exist & load correctly (no points awarded)
    # ---------------------------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File not found")
        return 0.0

    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded with {len(prs.slides)} slide(s)")
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        return 0.0

    # ---------------------------------------------------------------------
    # 1. Locate first table (Table 1)
    # ---------------------------------------------------------------------
    first_table = None
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_table:
                first_table = shape.table
                table_slide_index = slide_idx
                table_shape_index = shape_idx
                break
        if first_table is not None:
            break

    if first_table is None:
        print("✗ No table found in presentation – cannot verify task")
        return 0.0

    print(f"✓ Found first table on slide {table_slide_index + 1}, shape {table_shape_index + 1}")

    # ---------------------------------------------------------------------
    # 2. Requirement 1 – ALL borders removed (0.5 pts)
    # ---------------------------------------------------------------------
    borders_ok = True
    for row_idx, row in enumerate(first_table.rows):
        for col_idx, cell in enumerate(row.cells):
            tcPr = cell._tc.tcPr  # underlying XML element for table-cell properties
            # any child element whose localname starts with 'ln' represents a border
            for ln in tcPr.iterchildren():
                tag_local = etree.QName(ln).localname
                if tag_local.startswith('ln'):  # lnL, lnR, lnT, lnB, etc.
                    w_val = ln.get('w')  # border width in EMU (string)
                    try:
                        width_int = int(w_val) if w_val is not None else None
                    except ValueError:
                        width_int = None
                    # Any missing width or width > 0 ⇒ border present ⇒ fail
                    if width_int is None or width_int > 0:
                        print(f"   ✗ Border detected in cell ({row_idx},{col_idx}) with width={w_val}")
                        borders_ok = False
                        break
            if not borders_ok:
                break
        if not borders_ok:
            break

    if borders_ok:
        total_score += 0.5
        print("✓ All borders appear to be removed (0.5 pts)")
    else:
        print("✗ Borders not fully removed (0 pts)")

    # ---------------------------------------------------------------------
    # 3. Requirement 2 – Cell padding ≈ 0.2 cm (0.5 pts)
    # ---------------------------------------------------------------------
    TARGET_CM = 0.2
    TOLERANCE_CM = 0.1  # allow ±0.1 cm due to app/internal rounding
    EMU_PER_CM = 360000.0

    padding_ok = True
    for row_idx, row in enumerate(first_table.rows):
        for col_idx, cell in enumerate(row.cells):
            margins_emu = [cell.margin_left, cell.margin_top, cell.margin_right, cell.margin_bottom]
            for margin_emu in margins_emu:
                margin_cm = margin_emu / EMU_PER_CM
                if abs(margin_cm - TARGET_CM) > TOLERANCE_CM:
                    print(f"   ✗ Cell ({row_idx},{col_idx}) margin {margin_cm:.3f} cm outside tolerance")
                    padding_ok = False
                    break
            if not padding_ok:
                break
        if not padding_ok:
            break

    if padding_ok:
        total_score += 0.5
        print("✓ All cell paddings approximately 0.2 cm (0.5 pts)")
    else:
        print("✗ Cell padding not set correctly (0 pts)")

    # ---------------------------------------------------------------------
    # 4. Final score & report
    # ---------------------------------------------------------------------
    final_score = min(total_score, max_score)
    print(f"Total score: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score

# -----------------------------------------------------------------------------
# Self-test when run as a standalone script (uses the golden file path)
# -----------------------------------------------------------------------------
if __name__ == "__main__":
    GOLDEN_FILE = "/home/user/apply_no_borders_to_table_1_and_set_cell_padding_to_02_cm.pptx"
    verify_task(GOLDEN_FILE)

