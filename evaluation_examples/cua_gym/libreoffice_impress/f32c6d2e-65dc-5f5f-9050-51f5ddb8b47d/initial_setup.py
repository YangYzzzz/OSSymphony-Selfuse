"""
Initial Setup: Create a 6-slide Campus Survey presentation with slide 4 having
title 'Survey Results' and empty content area (no chart).
Task ID: impress_stu_025
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_025'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=None):
    """Helper to add a text box with formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Campus Life Survey 2025"
    slide1.placeholders[1].text = "Understanding Student Satisfaction & Engagement"

    # --- Slide 2: Introduction ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Introduction"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Survey Overview"
    body2.paragraphs[0].runs[0].font.bold = True
    body2.paragraphs[0].runs[0].font.size = Pt(20)

    p = body2.add_paragraph()
    p.text = "This survey was conducted across 4 university campuses during the Fall 2024 semester."
    p.level = 0

    p = body2.add_paragraph()
    p.text = "Total respondents: 2,847 undergraduate and graduate students"
    p.level = 0

    p = body2.add_paragraph()
    p.text = "Response rate: 68.3% of enrolled students"
    p.level = 0

    p = body2.add_paragraph()
    p.text = "Demographics: 52% female, 46% male, 2% non-binary"
    p.level = 0

    # --- Slide 3: Methodology ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Methodology"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Data Collection Methods"
    body3.paragraphs[0].runs[0].font.bold = True
    body3.paragraphs[0].runs[0].font.size = Pt(20)

    methods = [
        "Online questionnaire distributed via student email (Sept 15 - Oct 30, 2024)",
        "5-point Likert scale used for satisfaction metrics",
        "Focus group sessions with 120 randomly selected participants",
        "Confidence interval: 95% with margin of error +/- 1.8%",
        "Data validated against enrollment records for accuracy",
    ]
    for m in methods:
        p = body3.add_paragraph()
        p.text = m
        p.level = 0

    # --- Slide 4: Survey Results (EMPTY - no chart) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Survey Results"
    # Leave the content placeholder with minimal text to indicate empty area
    body4 = slide4.placeholders[1].text_frame
    body4.text = ""  # empty content area

    # --- Slide 5: Key Findings ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Key Findings"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Highlights from the Survey"
    body5.paragraphs[0].runs[0].font.bold = True
    body5.paragraphs[0].runs[0].font.size = Pt(20)

    findings = [
        "63% of students expressed overall satisfaction with campus life",
        "Library resources rated highest at 4.2/5.0 average score",
        "Campus dining received the lowest satisfaction at 3.1/5.0",
        "Housing quality showed significant variation across campuses",
        "Mental health services demand increased 34% year-over-year",
    ]
    for f in findings:
        p = body5.add_paragraph()
        p.text = f
        p.level = 0

    # --- Slide 6: Conclusion & Next Steps ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Conclusion & Next Steps"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Recommended Actions"
    body6.paragraphs[0].runs[0].font.bold = True
    body6.paragraphs[0].runs[0].font.size = Pt(20)

    actions = [
        "Expand mental health counseling hours by 40% starting Spring 2025",
        "Conduct dining services vendor review with student advisory panel",
        "Invest $2.3M in dormitory renovations for East and West campuses",
        "Launch monthly student feedback portal for real-time tracking",
        "Schedule follow-up survey for May 2025 to measure improvement",
    ]
    for a in actions:
        p = body6.add_paragraph()
        p.text = a
        p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
