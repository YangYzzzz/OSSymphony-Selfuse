"""
Reward Script: Create a pie chart on slide 4 with survey results
Task ID: impress_stu_025
Domain: libreoffice_impress
Scoring:
  Component 1: Chart exists on slide 4 (0.20)
  Component 2: Chart type is PIE (0.15)
  Component 3: Chart title matches (0.20)
  Component 4: Categories match (0.25)
  Component 5: Values match (0.20)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_025'

# Expected data
EXPECTED_CATEGORIES = ['Strongly Agree', 'Agree', 'Neutral', 'Disagree', 'Strongly Disagree']
EXPECTED_VALUES = [35.0, 28.0, 20.0, 12.0, 5.0]
EXPECTED_TITLE = 'Student Satisfaction Survey Results'


def persist_app_state():
    """Save any open LibreOffice documents before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)

    # Find chart shape on slide 4
    chart_shape = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            chart_shape = shape
            break

    # Component 1: Chart exists on slide 4 (0.20 points)
    try:
        if chart_shape is not None:
            print(f"PASS: Component 1 -- Chart found on slide 4 (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 -- No chart found on slide 4")
            # Without a chart, no further checks possible
            final_score = min(total_score, 1.0)
            print(f"\nScore: {total_score}/1.0")
            print(f"REWARD: {final_score}")
            return final_score
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    chart = chart_shape.chart

    # Component 2: Chart type is PIE (0.15 points)
    try:
        # PIE chart type enum value is 5
        chart_type = chart.chart_type
        if chart_type == 5:  # PIE
            print(f"PASS: Component 2 -- Chart type is PIE (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 -- Expected PIE chart (type 5), found type {chart_type}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Chart title is 'Student Satisfaction Survey Results' (0.20 points)
    try:
        if chart.has_title:
            actual_title = chart.chart_title.text_frame.text.strip()
            if actual_title == EXPECTED_TITLE:
                print(f"PASS: Component 3 -- Chart title matches exactly (0.20 pts)")
                total_score += 0.20
            elif actual_title.lower() == EXPECTED_TITLE.lower():
                # Accept case-insensitive match with partial credit
                print(f"PARTIAL: Component 3 -- Chart title matches case-insensitively: '{actual_title}' (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 3 -- Expected title '{EXPECTED_TITLE}', found '{actual_title}'")
        else:
            print(f"FAIL: Component 3 -- Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Categories match exactly (0.25 points)
    try:
        plot = chart.plots[0]
        actual_cats = list(plot.categories)

        if actual_cats == EXPECTED_CATEGORIES:
            print(f"PASS: Component 4 -- All 5 categories match exactly (0.25 pts)")
            total_score += 0.25
        else:
            # Check partial match: count how many categories are correct
            matching = 0
            for exp_cat in EXPECTED_CATEGORIES:
                if exp_cat in actual_cats:
                    matching += 1

            if matching == len(EXPECTED_CATEGORIES):
                # All categories present but possibly in wrong order
                print(f"PARTIAL: Component 4 -- All categories present but order differs (0.15 pts)")
                total_score += 0.15
            elif matching > 0:
                partial = 0.25 * (matching / len(EXPECTED_CATEGORIES)) * 0.5
                print(f"PARTIAL: Component 4 -- {matching}/{len(EXPECTED_CATEGORIES)} categories found (actual: {actual_cats}) ({partial:.2f} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 4 -- Expected categories {EXPECTED_CATEGORIES}, found {actual_cats}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Values match (0.20 points)
    try:
        series_values = list(chart.series[0].values)

        if len(series_values) == len(EXPECTED_VALUES):
            # Check if values match (allow small floating point tolerance)
            mismatches = [
                (actual, expected)
                for actual, expected in zip(series_values, EXPECTED_VALUES)
                if abs(float(actual) - float(expected)) > 0.5
            ]

            if len(mismatches) == 0:
                print(f"PASS: Component 5 -- All values match: {series_values} (0.20 pts)")
                total_score += 0.20
            else:
                # Check if values match but in different order (paired with categories)
                actual_pairs = sorted(zip(list(plot.categories), series_values))
                expected_pairs = sorted(zip(EXPECTED_CATEGORIES, EXPECTED_VALUES))
                reordered_match = all(
                    a[0] == e[0] and abs(float(a[1]) - float(e[1])) < 0.5
                    for a, e in zip(actual_pairs, expected_pairs)
                )
                if reordered_match:
                    print(f"PARTIAL: Component 5 -- Values match when reordered by category (0.15 pts)")
                    total_score += 0.15
                else:
                    print(f"FAIL: Component 5 -- Expected values {EXPECTED_VALUES}, found {series_values}")
        else:
            print(f"FAIL: Component 5 -- Expected {len(EXPECTED_VALUES)} values, found {len(series_values)}: {series_values}")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
