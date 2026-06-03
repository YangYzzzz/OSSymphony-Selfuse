"""
Initial Setup: Create a 10-slide presentation with images and text for PDF export task
Task ID: impress_el_009
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import struct
import zlib

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_el_009'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_minimal_png(width, height, r, g, b, filename):
    """Create a minimal solid-color PNG file."""

    def make_chunk(chunk_type, data):
        chunk = chunk_type + data
        crc = struct.pack('>I', zlib.crc32(chunk) & 0xFFFFFFFF)
        return struct.pack('>I', len(data)) + chunk + crc

    sig = b'\x89PNG\r\n\x1a\n'
    ihdr_data = struct.pack('>IIBBBBB', width, height, 8, 2, 0, 0, 0)
    raw = b''
    row = bytes([r, g, b] * width)
    for _ in range(height):
        raw += b'\x00' + row
    idat_data = zlib.compress(raw)

    png_bytes = sig
    png_bytes += make_chunk(b'IHDR', ihdr_data)
    png_bytes += make_chunk(b'IDAT', idat_data)
    png_bytes += make_chunk(b'IEND', b'')

    with open(filename, 'wb') as f:
        f.write(png_bytes)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Create placeholder images
    img_configs = [
        (400, 300, 41, 128, 185, f'{WORKDIR}/chart_revenue.png'),
        (400, 300, 46, 139, 87, f'{WORKDIR}/chart_growth.png'),
        (400, 300, 230, 126, 34, f'{WORKDIR}/chart_market.png'),
        (400, 300, 142, 68, 173, f'{WORKDIR}/chart_satisfaction.png'),
        (400, 300, 192, 57, 43, f'{WORKDIR}/chart_pipeline.png'),
        (400, 300, 44, 62, 80, f'{WORKDIR}/chart_budget.png'),
        (400, 300, 39, 174, 96, f'{WORKDIR}/chart_timeline.png'),
        (400, 300, 211, 84, 0, f'{WORKDIR}/chart_risk.png'),
        (400, 300, 52, 73, 94, f'{WORKDIR}/chart_resources.png'),
        (400, 300, 22, 160, 133, f'{WORKDIR}/chart_kpi.png'),
    ]
    for w, h, r, g, b, path in img_configs:
        create_minimal_png(w, h, r, g, b, path)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q4 2025 Accessibility & Compliance Report"
    slide1.placeholders[1].text = "Global Operations Division — Prepared by Ananya Patel"
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x2A, 0x4A)
    for shape in slide1.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    title2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = title2.text_frame
    p = tf.paragraphs[0]
    p.text = "Executive Summary"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

    body2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(7), Inches(4.5))
    tf2 = body2.text_frame
    tf2.word_wrap = True
    bullets = [
        "Overall accessibility compliance improved to 94.2% across all digital platforms",
        "WCAG 2.1 Level AA conformance achieved for 87% of customer-facing applications",
        "Screen reader compatibility testing expanded to cover 12 additional product modules",
        "Remediation backlog reduced by 38% compared to Q3 2025 through dedicated sprint cycles",
        "Accessibility training completion rate reached 91% among development teams globally",
    ]
    for i, text in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = text
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    slide2.shapes.add_picture(img_configs[0][5], Inches(8.5), Inches(1.8), Inches(4), Inches(3))

    # --- Slide 3: Revenue Impact ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    title3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = title3.text_frame
    p = tf.paragraphs[0]
    p.text = "Revenue Impact of Accessibility Initiatives"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

    body3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.5))
    tf3 = body3.text_frame
    tf3.word_wrap = True
    lines = [
        "Accessible product lines generated $4.8M in additional revenue this quarter",
        "Customer retention improved by 12% among users requiring assistive technologies",
        "New market penetration in government sector yielded $2.1M in contracts",
        "Cost avoidance from proactive compliance: estimated $890K in potential litigation savings",
    ]
    for i, text in enumerate(lines):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = text
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.size = Pt(16)

    slide3.shapes.add_picture(img_configs[1][5], Inches(8), Inches(1.8), Inches(4.5), Inches(3.5))

    # --- Slide 4: Compliance Metrics ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    title4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = title4.text_frame
    p = tf.paragraphs[0]
    p.text = "WCAG 2.1 Compliance Metrics by Product"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

    table_shape = slide4.shapes.add_table(6, 4, Inches(0.8), Inches(1.8), Inches(8), Inches(3.5))
    table = table_shape.table
    headers = ["Product", "Level A", "Level AA", "Level AAA"]
    data = [
        ["Customer Portal", "100%", "96%", "72%"],
        ["Mobile Banking App", "100%", "91%", "65%"],
        ["Internal HR Platform", "98%", "88%", "58%"],
        ["E-Commerce Storefront", "100%", "94%", "70%"],
        ["Document Management", "97%", "85%", "52%"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    slide4.shapes.add_picture(img_configs[2][5], Inches(9.5), Inches(1.8), Inches(3), Inches(2.5))

    # --- Slide 5: Screen Reader Testing ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    title5 = slide5.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = title5.text_frame
    p = tf.paragraphs[0]
    p.text = "Screen Reader Compatibility Results"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

    body5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(7), Inches(4.5))
    tf5 = body5.text_frame
    tf5.word_wrap = True
    sr_items = [
        "JAWS 2025: 97.3% navigation success rate across all tested interfaces",
        "NVDA 2025.1: 95.8% compatibility score with dynamic content updates",
        "VoiceOver (macOS 15): 94.1% landmark detection accuracy",
        "TalkBack (Android 16): 92.6% gesture navigation reliability",
        "Narrator (Windows 12): 96.2% form interaction success rate",
    ]
    for i, text in enumerate(sr_items):
        if i == 0:
            p = tf5.paragraphs[0]
        else:
            p = tf5.add_paragraph()
        p.text = text
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.size = Pt(16)

    slide5.shapes.add_picture(img_configs[3][5], Inches(8.5), Inches(1.8), Inches(4), Inches(3))

    # --- Slide 6: Remediation Progress ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    title6 = slide6.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = title6.text_frame
    p = tf.paragraphs[0]
    p.text = "Remediation Backlog Progress"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

    body6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6), Inches(4.5))
    tf6 = body6.text_frame
    tf6.word_wrap = True
    rem_items = [
        "Critical issues resolved: 142 of 148 (96%)",
        "High priority items completed: 287 of 310 (93%)",
        "Medium priority fixes deployed: 534 of 620 (86%)",
        "Low priority improvements scheduled: 891 items for Q1 2026",
        "Average resolution time decreased from 14 days to 8.5 days",
    ]
    for i, text in enumerate(rem_items):
        if i == 0:
            p = tf6.paragraphs[0]
        else:
            p = tf6.add_paragraph()
        p.text = text
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.size = Pt(16)

    slide6.shapes.add_picture(img_configs[4][5], Inches(8), Inches(2), Inches(4.5), Inches(3))

    # --- Slide 7: Training & Awareness ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    title7 = slide7.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = title7.text_frame
    p = tf.paragraphs[0]
    p.text = "Training & Awareness Programs"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

    table7 = slide7.shapes.add_table(5, 3, Inches(0.8), Inches(1.8), Inches(7), Inches(3))
    t7 = table7.table
    t7_headers = ["Program", "Participants", "Completion Rate"]
    t7_data = [
        ["Intro to Web Accessibility", "1,245", "94%"],
        ["Advanced ARIA Techniques", "486", "88%"],
        ["Mobile A11y Best Practices", "732", "91%"],
        ["Document Accessibility", "958", "89%"],
    ]
    for c, h in enumerate(t7_headers):
        cell = t7.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row_data in enumerate(t7_data, 1):
        for c, val in enumerate(row_data):
            t7.cell(r, c).text = val

    slide7.shapes.add_picture(img_configs[5][5], Inches(8.5), Inches(1.8), Inches(4), Inches(3))

    # --- Slide 8: Budget Allocation ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    title8 = slide8.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = title8.text_frame
    p = tf.paragraphs[0]
    p.text = "Budget Allocation & Spending"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

    body8 = slide8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.5))
    tf8 = body8.text_frame
    tf8.word_wrap = True
    budget_items = [
        "Total Q4 budget: $3.2M | Spent: $2.87M (89.7% utilization)",
        "Automated testing tools: $420K (Axe Enterprise, Pa11y CI, WAVE API)",
        "Third-party audits: $380K (Deque, Level Access, WebAIM)",
        "Training program development: $290K across 4 new curricula",
        "Assistive technology procurement: $185K for internal testing lab",
        "Remaining funds: $325K carried forward to Q1 2026 remediation sprint",
    ]
    for i, text in enumerate(budget_items):
        if i == 0:
            p = tf8.paragraphs[0]
        else:
            p = tf8.add_paragraph()
        p.text = text
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.size = Pt(15)

    slide8.shapes.add_picture(img_configs[6][5], Inches(8), Inches(2), Inches(4.5), Inches(3.5))

    # --- Slide 9: Risk Assessment ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    title9 = slide9.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = title9.text_frame
    p = tf.paragraphs[0]
    p.text = "Risk Assessment & Mitigation"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

    body9 = slide9.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.5))
    tf9 = body9.text_frame
    tf9.word_wrap = True
    risks = [
        "HIGH: European Accessibility Act enforcement begins June 2025 — 3 products need updates",
        "MEDIUM: Section 508 refresh audit scheduled for March 2026 — preparation in progress",
        "LOW: WCAG 3.0 draft may introduce new conformance model — monitoring working group",
        "Mitigation: Dedicated compliance sprint planned for Jan-Feb 2026 with 8 FTE allocation",
    ]
    for i, text in enumerate(risks):
        if i == 0:
            p = tf9.paragraphs[0]
        else:
            p = tf9.add_paragraph()
        p.text = text
        p.space_after = Pt(10)
        for run in p.runs:
            run.font.size = Pt(15)

    slide9.shapes.add_picture(img_configs[7][5], Inches(8.5), Inches(1.8), Inches(4), Inches(3))

    # --- Slide 10: Next Steps & Recommendations ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    title10 = slide10.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = title10.text_frame
    p = tf.paragraphs[0]
    p.text = "Next Steps & Recommendations"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

    body10 = slide10.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(7), Inches(4.5))
    tf10 = body10.text_frame
    tf10.word_wrap = True
    steps = [
        "1. Complete EAA compliance gap analysis for all EU-facing products by Jan 31, 2026",
        "2. Deploy AI-powered accessibility testing in CI/CD pipeline for 5 remaining products",
        "3. Establish accessibility champion network: 2 representatives per engineering team",
        "4. Launch quarterly external user testing with disability advocacy organizations",
        "5. Increase automated test coverage from 68% to 85% of WCAG 2.1 success criteria",
        "6. Submit annual VPAT updates for all commercial products by February 2026",
    ]
    for i, text in enumerate(steps):
        if i == 0:
            p = tf10.paragraphs[0]
        else:
            p = tf10.add_paragraph()
        p.text = text
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.size = Pt(16)

    slide10.shapes.add_picture(img_configs[8][5], Inches(8.5), Inches(2), Inches(4), Inches(3))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT} with {len(prs.slides)} slides')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=3.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
