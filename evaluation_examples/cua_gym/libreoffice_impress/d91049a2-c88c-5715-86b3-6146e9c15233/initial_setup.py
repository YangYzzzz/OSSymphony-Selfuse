"""
Initial Setup: Export slide 5 as PNG to Desktop
Task ID: impress_gf3_003
Domain: libreoffice_impress

Creates a 12-slide Design_Portfolio.pptx with slide 5 having a full-bleed
image background, two overlapping text boxes, and a shape. Opens in Impress.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFilter

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_003'
FILENAME = 'Design_Portfolio.pptx'
OUTPUT = f'{WORKDIR}/{FILENAME}'
BG_IMAGE_PATH = f'{WORKDIR}/_slide5_bg.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_background_image():
    """Create an abstract gradient background image for slide 5."""
    width, height = 1920, 1440
    img = Image.new('RGB', (width, height))
    draw = ImageDraw.Draw(img)

    # Create a gradient with geometric shapes
    for y in range(height):
        r = int(20 + (y / height) * 60)
        g = int(40 + (y / height) * 80)
        b = int(100 + (y / height) * 100)
        draw.line([(0, y), (width, y)], fill=(r, g, b))

    # Add some geometric accents
    draw.ellipse([100, 200, 600, 700], fill=(60, 120, 180), outline=(80, 140, 200))
    draw.rectangle([1200, 300, 1700, 900], fill=(40, 90, 150), outline=(60, 110, 170))
    draw.polygon([(900, 100), (1100, 500), (700, 500)], fill=(50, 100, 160))

    # Slight blur for smoothness
    img = img.filter(ImageFilter.GaussianBlur(radius=3))
    img.save(BG_IMAGE_PATH)
    return BG_IMAGE_PATH


def create_presentation():
    prs = Presentation()
    # Standard 13.33 x 7.5 inch (33.87 x 19.05 cm) widescreen
    # Actually python-pptx default is 10x7.5 (standard 4:3)
    # Use default which is 10x7.5 inches = 25.4 x 19.05 cm
    # At 200 DPI: 2000 x 1500 pixels
    # For 33.87 x 25.4 cm = 13.33 x 10 inches at 200 DPI = 2667 x 2000
    # Task says "approximately 2540x1905" which is 12.7 x 9.525 inches at 200 DPI
    # That corresponds to standard slide 33.867cm x 19.05cm widescreen (13.333 x 7.5)
    # Wait: 2540/200 = 12.7 inches = 32.258 cm, 1905/200 = 9.525 inches = 24.19 cm
    # Actually the task says "33.87x25.4cm" which is standard A4-ish
    # 33.867cm = 13.333in, 25.4cm = 10in => 200 DPI => 2667x2000
    # But task says ~2540x1905. Let me check: default pptx is 10x7.5 => 2000x1500
    # Widescreen 13.333x7.5 => 2667x1500
    # Hmm, 2540x1905: 2540/200=12.7in=32.258cm, 1905/200=9.525in=24.19cm
    # That's close to 25.4cm height... Let me just use default and the task says "approximately"
    # Standard slide is 10x7.5 inches. At 200 DPI = 2000x1500.
    # But task mentions 2540x1905 and "33.87x25.4cm". Let me set custom size.
    # 25.4cm = 10 inches width, 19.05cm = 7.5 inches height => standard
    # 33.867cm width = 13.333 inches. Hmm.
    # 2540 pixels / 200 DPI = 12.7 inches. 1905 / 200 = 9.525 inches.
    # Actually: 25.4cm = 10in (width), 19.05cm = 7.5in (height) is standard 4:3
    # At 200 DPI that's 2000x1500. Not 2540x1905.
    # Maybe they mean: width=33.867cm=13.333in, height=25.4cm=10in => portrait?
    # 13.333*200=2667, 10*200=2000. Still not matching.
    # Let me just go with: slide 25.4cm x 19.05cm (10x7.5in standard)
    # The golden patch will handle the actual export dimensions.
    # For 2540x1905: that's exactly 25.4cm x 19.05cm at 100 pixels/cm = 254 DPI equivalent
    # Or: in metric at 200 DPI: 25.4cm * (200/2.54) = 2000px. Not 2540.
    # Actually 2540 = 25.4 * 100. So 100 px/cm = 254 DPI (close to 200 with rounding)
    # Let me not overthink this. The golden_patch needs to export using soffice which
    # will determine the actual resolution. I'll use standard slide dimensions.

    slide_width = prs.slide_width   # default 10 inches
    slide_height = prs.slide_height  # default 7.5 inches

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Meridian Design Studio"
    slide1.placeholders[1].text = "Portfolio & Creative Showcase 2025"

    # --- Slide 2: About Us ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "About Our Studio"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Founded in 2018 by Elena Vasquez and Takeshi Nakamura"
    p2 = body2.add_paragraph()
    p2.text = "We specialize in brand identity, digital experiences, and spatial design"
    p3 = body2.add_paragraph()
    p3.text = "Our team of 14 designers has delivered over 200 projects globally"

    # --- Slide 3: Services Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Our Services"
    body3 = slide3.placeholders[1].text_frame
    services = [
        "Brand Identity & Visual Systems",
        "UI/UX Design & Prototyping",
        "Motion Graphics & Animation",
        "Environmental & Spatial Design",
        "Print Design & Editorial Layout",
    ]
    body3.text = services[0]
    for svc in services[1:]:
        p = body3.add_paragraph()
        p.text = svc

    # --- Slide 4: Client Logos ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title_box = slide4.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_box.text_frame.paragraphs[0].text = "Trusted By Leading Brands"
    title_box.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
    title_box.text_frame.paragraphs[0].runs[0].font.bold = True

    clients = ["Aurelia Cosmetics", "NovaTech Industries", "Greenfield Agriculture",
               "Sapphire Hotels", "Velocity Sports", "Cascade Media Group"]
    y_pos = Inches(2)
    for i, client in enumerate(clients):
        col = i % 3
        row = i // 3
        tb = slide4.shapes.add_textbox(
            Inches(0.5 + col * 3.2), y_pos + Inches(row * 2),
            Inches(2.8), Inches(1.5)
        )
        tf = tb.text_frame
        tf.paragraphs[0].text = client
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].runs[0].font.size = Pt(16)
        tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 5: Featured Project (full-bleed bg + overlapping text + shape) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    # Full-bleed background image
    bg_path = create_background_image()
    slide5.shapes.add_picture(
        bg_path, 0, 0, slide_width, slide_height
    )

    # Overlapping text box 1 (project title)
    tx1 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(6), Inches(2))
    tf1 = tx1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "Project Aurora"
    p1.runs[0].font.size = Pt(44)
    p1.runs[0].font.bold = True
    p1.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p1b = tf1.add_paragraph()
    p1b.text = "Brand Identity Redesign for Aurelia Cosmetics"
    run1b = p1b.runs[0]
    run1b.font.size = Pt(18)
    run1b.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

    # Overlapping text box 2 (project details, overlaps with box 1)
    tx2 = slide5.shapes.add_textbox(Inches(1.2), Inches(3.0), Inches(5.5), Inches(2.5))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2a = tf2.paragraphs[0]
    p2a.text = "Scope: Visual identity, packaging, digital presence"
    p2a.runs[0].font.size = Pt(14)
    p2a.runs[0].font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    p2b = tf2.add_paragraph()
    p2b.text = "Duration: 8 weeks | Team: 4 designers"
    p2b.runs[0].font.size = Pt(14)
    p2b.runs[0].font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    p2c = tf2.add_paragraph()
    p2c.text = "Result: 340% increase in brand recognition"
    p2c.runs[0].font.size = Pt(14)
    p2c.runs[0].font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)

    # Decorative shape (rounded rectangle accent)
    shape5 = slide5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7), Inches(5), Inches(2.5), Inches(1.5)
    )
    shape5.fill.solid()
    shape5.fill.fore_color.rgb = RGBColor(0xFF, 0xA5, 0x00)
    shape5.line.color.rgb = RGBColor(0xFF, 0xC0, 0x40)
    tf_s = shape5.text_frame
    tf_s.paragraphs[0].text = "Award Winner"
    tf_s.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_s.paragraphs[0].runs[0].font.size = Pt(16)
    tf_s.paragraphs[0].runs[0].font.bold = True
    tf_s.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # --- Slide 6: Project Details ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Aurora - Design Process"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Phase 1: Research & Discovery (2 weeks)"
    for phase in ["Phase 2: Concept Development (2 weeks)",
                  "Phase 3: Design Refinement (2 weeks)",
                  "Phase 4: Production & Delivery (2 weeks)"]:
        p = body6.add_paragraph()
        p.text = phase

    # --- Slide 7: Another Project ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    fill7 = slide7.background.fill
    fill7.solid()
    fill7.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    tx7 = slide7.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tf7 = tx7.text_frame
    p7 = tf7.paragraphs[0]
    p7.text = "Project Zenith"
    p7.runs[0].font.size = Pt(40)
    p7.runs[0].font.bold = True
    p7.runs[0].font.color.rgb = RGBColor(0xEC, 0xF0, 0xF1)
    p7b = tf7.add_paragraph()
    p7b.text = "NovaTech Industries - Product Launch Campaign"
    p7b.runs[0].font.size = Pt(20)
    p7b.runs[0].font.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)

    # --- Slide 8: Statistics ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title8.text_frame.paragraphs[0].text = "Impact & Results"
    title8.text_frame.paragraphs[0].runs[0].font.size = Pt(30)
    title8.text_frame.paragraphs[0].runs[0].font.bold = True

    stats = [("200+", "Projects Delivered"), ("14", "Team Members"),
             ("45", "Industry Awards"), ("98%", "Client Satisfaction")]
    for i, (num, label) in enumerate(stats):
        x = Inches(0.5 + i * 2.4)
        tb_num = slide8.shapes.add_textbox(x, Inches(2.5), Inches(2), Inches(1.5))
        tb_num.text_frame.paragraphs[0].text = num
        tb_num.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tb_num.text_frame.paragraphs[0].runs[0].font.size = Pt(48)
        tb_num.text_frame.paragraphs[0].runs[0].font.bold = True
        tb_num.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x27, 0x72, 0xDB)

        tb_lbl = slide8.shapes.add_textbox(x, Inches(4), Inches(2), Inches(1))
        tb_lbl.text_frame.paragraphs[0].text = label
        tb_lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tb_lbl.text_frame.paragraphs[0].runs[0].font.size = Pt(14)

    # --- Slide 9: Testimonial ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    fill9 = slide9.background.fill
    fill9.solid()
    fill9.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF0)

    quote_box = slide9.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(4))
    qtf = quote_box.text_frame
    qtf.word_wrap = True
    qp = qtf.paragraphs[0]
    qp.text = '"Meridian transformed our entire brand presence. Their attention to detail and creative vision exceeded every expectation."'
    qp.alignment = PP_ALIGN.CENTER
    qp.runs[0].font.size = Pt(22)
    qp.runs[0].font.italic = True
    qp.runs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    attr = qtf.add_paragraph()
    attr.text = "- Camille Beaumont, CEO of Aurelia Cosmetics"
    attr.alignment = PP_ALIGN.CENTER
    attr.runs[0].font.size = Pt(14)
    attr.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # --- Slide 10: Process Overview ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "Our Creative Process"
    body10 = slide10.placeholders[1].text_frame
    steps = ["Discover: Understanding your vision and goals",
             "Define: Establishing strategy and creative direction",
             "Design: Crafting visual solutions iteratively",
             "Deliver: Polished assets ready for production"]
    body10.text = steps[0]
    for s in steps[1:]:
        p = body10.add_paragraph()
        p.text = s

    # --- Slide 11: Pricing ---
    slide11 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    t11 = slide11.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    t11.text_frame.paragraphs[0].text = "Engagement Models"
    t11.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
    t11.text_frame.paragraphs[0].runs[0].font.bold = True

    tiers = [("Essentials", "$5,000 - $15,000", "Logo, color palette, typography"),
             ("Professional", "$15,000 - $40,000", "Full brand identity system"),
             ("Enterprise", "$40,000+", "Complete brand transformation")]
    for i, (name, price, desc) in enumerate(tiers):
        x = Inches(0.3 + i * 3.2)
        box = slide11.shapes.add_textbox(x, Inches(2), Inches(3), Inches(4))
        btf = box.text_frame
        btf.word_wrap = True
        bp = btf.paragraphs[0]
        bp.text = name
        bp.runs[0].font.size = Pt(22)
        bp.runs[0].font.bold = True
        bp.runs[0].font.color.rgb = RGBColor(0x27, 0x72, 0xDB)
        bp2 = btf.add_paragraph()
        bp2.text = price
        bp2.runs[0].font.size = Pt(18)
        bp2.runs[0].font.bold = True
        bp3 = btf.add_paragraph()
        bp3.text = desc
        bp3.runs[0].font.size = Pt(13)
        bp3.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # --- Slide 12: Contact ---
    slide12 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    fill12 = slide12.background.fill
    fill12.solid()
    fill12.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    ct = slide12.shapes.add_textbox(Inches(2), Inches(1.5), Inches(6), Inches(5))
    ctf = ct.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.text = "Let's Create Something Extraordinary"
    cp.alignment = PP_ALIGN.CENTER
    cp.runs[0].font.size = Pt(32)
    cp.runs[0].font.bold = True
    cp.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    for line in ["", "hello@meridiandesign.studio",
                 "+1 (415) 555-0192", "meridiandesign.studio",
                 "", "San Francisco | Tokyo | Berlin"]:
        lp = ctf.add_paragraph()
        lp.text = line
        lp.alignment = PP_ALIGN.CENTER
        if lp.runs:
            lp.runs[0].font.size = Pt(16)
            lp.runs[0].font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Slide count: {len(prs.slides)}')

    # Clean up temp image
    if os.path.exists(BG_IMAGE_PATH):
        os.remove(BG_IMAGE_PATH)

    # Ensure Desktop exists and has NO exported PNG
    os.makedirs(f'{WORKDIR}/Desktop', exist_ok=True)
    export_path = f'{WORKDIR}/Desktop/slide5_export.png'
    if os.path.exists(export_path):
        os.remove(export_path)

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_presentation()
