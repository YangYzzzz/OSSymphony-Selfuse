"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 273 is the pricing summary and I need a quick table there: exactly 2 rows and 5 columns. After it’s in place, please right-align the text in the header row so the labels line up neatly on the right edge.
Generated: 2025-09-10 20:43:19
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation
from pptx.enum.text import PP_ALIGN


def verify_pricing_table(file_path: str) -> float:
    """Verify that slide 273 contains a 2x5 table and that all header cells
    (first row) are right-aligned. Returns a progressive score between 0.0
    and 1.0.
    """
    print(f"Verifying file: {file_path}")

    score = 0.0          # progressive score
    max_score = 1.0      # cap

    # ------------------------------------------------------------------
    # 1) Load the presentation (prerequisite – no points awarded)
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File not found – task failed")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded successfully with {len(prs.slides)} slides (0 pts)")
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # ------------------------------------------------------------------
    # 2) Locate slide 273 (zero-based index 272)
    # ------------------------------------------------------------------
    target_index = 272  # Slide 273 in zero-based numbering
    if target_index >= len(prs.slides):
        print(f"✗ Slide 273 not present – only {len(prs.slides)} slides found")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[target_index]
    print("✓ Slide 273 located (0 pts)")

    # ------------------------------------------------------------------
    # 3) Look for a table with exactly 2 rows and 5 columns
    # ------------------------------------------------------------------
    desired_table = None
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            tbl = shape.table
            print(f"   Found table: {len(tbl.rows)} rows × {len(tbl.columns)} columns")
            if len(tbl.rows) == 2 and len(tbl.columns) == 5:
                desired_table = tbl
                break

    if desired_table is None:
        print("✗ No 2×5 table found on slide 273")
    else:
        print("✓ Correctly-sized 2×5 table found (+0.5 pts)")
        score += 0.5

        # ------------------------------------------------------------------
        # 4) Verify header row (row 0) is right-aligned in every column
        # ------------------------------------------------------------------
        header_ok = True
        for col_idx in range(len(desired_table.columns)):
            cell = desired_table.cell(0, col_idx)
            # If the cell has multiple paragraphs, check the first one (common)
            if not cell.text_frame.paragraphs:
                print(f"   ✗ Header cell {col_idx} contains no paragraph")
                header_ok = False
                continue

            alignment = cell.text_frame.paragraphs[0].alignment
            print(f"   Header cell {col_idx} alignment: {alignment}")
            if alignment != PP_ALIGN.RIGHT:
                header_ok = False

        if header_ok:
            print("✓ All header cells are right-aligned (+0.5 pts)")
            score += 0.5
        else:
            print("✗ One or more header cells are not right-aligned")

    # ------------------------------------------------------------------
    # 5) Final scoring
    # ------------------------------------------------------------------
    final_score = min(score, max_score)
    print(f"Final score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score


# ----------------------------------------------------------------------
# Execute verification when run as a script
# ----------------------------------------------------------------------
if __name__ == "__main__":
    verify_pricing_table("/home/user/slide_273_is_the_pricing_summary_and_i_need_a_quick_table_there_exactly_2_rows_and_5_columns_after_i_golden.pptx")
