"""
Reward Script: Rotate header row text to 90 degrees vertical in table on slide 3
Task ID: impress_teach_089
Domain: libreoffice_impress
Scoring:
  Component 1 (0.6): Header cells have vert="vert" attribute (progressive by count)
  Component 2 (0.4): All headers rotated + text preserved + non-header rows unrotated
"""

import os
import re
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_089'
EXPECTED_HEADERS = [
    'Regional Territory Name',
    'Total Revenue Generated (USD)',
    'Year-over-Year Growth Rate (%)',
    'Customer Acquisition Count',
    'Average Transaction Value (USD)',
    'Net Promoter Score (NPS)',
]
NUM_HEADER_COLS = 6


def persist_app_state(domain):
    """Try to save any unsaved LibreOffice changes."""
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print(f"PERSIST: ctrl+s sent for {domain}")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: find table on slide 3
    slide3 = prs.slides[2]
    table = None
    for shape in slide3.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            break

    if table is None:
        print("FAIL: No table found on slide 3")
        print("REWARD: 0.0")
        return 0.0

    num_cols = len(table.columns)
    if num_cols < NUM_HEADER_COLS:
        print(f"FAIL: Table has {num_cols} columns, expected {NUM_HEADER_COLS}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Header cells have vert="vert" attribute (0.6 points, progressive)
    # "vert" in OOXML bodyPr means 90-degree vertical text rotation
    try:
        rotated_count = 0
        for c in range(NUM_HEADER_COLS):
            cell = table.cell(0, c)
            tc_xml = cell._tc.xml
            vert_match = re.search(r'vert="([^"]+)"', tc_xml)
            if vert_match and vert_match.group(1) == 'vert':
                rotated_count += 1
                print(f"  PASS: Cell(0,{c}) has vert='vert'")
            else:
                vert_val = vert_match.group(1) if vert_match else 'NOT SET'
                print(f"  FAIL: Cell(0,{c}) vert='{vert_val}', expected 'vert'")

        if rotated_count > 0:
            component1_score = 0.6 * (rotated_count / NUM_HEADER_COLS)
            total_score += component1_score
            print(f"PASS: Component 1 — {rotated_count}/{NUM_HEADER_COLS} header cells rotated ({component1_score:.2f} pts)")
        else:
            print(f"FAIL: Component 1 — No header cells have vertical rotation")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: All headers rotated + text preserved + non-header rows unrotated (0.4 points)
    # This is a compound check anchored to the rotation change
    try:
        all_rotated = (rotated_count == NUM_HEADER_COLS)

        # Check header text is preserved
        text_mismatches = 0
        for c in range(NUM_HEADER_COLS):
            actual_text = table.cell(0, c).text.strip()
            expected_text = EXPECTED_HEADERS[c]
            if actual_text != expected_text:
                print(f"  FAIL: Header text mismatch col {c}: '{actual_text}' != '{expected_text}'")
                text_mismatches += 1

        # Check non-header rows are NOT rotated (data rows should remain horizontal)
        non_header_issues = 0
        num_rows = len(table.rows)
        for r in range(1, min(num_rows, 3)):  # spot-check rows 1-2
            for c in range(NUM_HEADER_COLS):
                cell = table.cell(r, c)
                tc_xml = cell._tc.xml
                vert_match = re.search(r'vert="([^"]+)"', tc_xml)
                if vert_match and vert_match.group(1) != 'horz':
                    print(f"  WARN: Non-header Cell({r},{c}) unexpectedly has vert='{vert_match.group(1)}'")
                    non_header_issues += 1

        if all_rotated and text_mismatches == 0 and non_header_issues == 0:
            print(f"PASS: Component 2 — All headers rotated, text preserved, non-headers horizontal (0.4 pts)")
            total_score += 0.4
        else:
            reasons = []
            if not all_rotated:
                reasons.append(f"only {rotated_count}/{NUM_HEADER_COLS} rotated")
            if text_mismatches > 0:
                reasons.append("header text modified")
            if non_header_issues > 0:
                reasons.append("non-header rows also rotated")
            print(f"FAIL: Component 2 — {', '.join(reasons)}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
