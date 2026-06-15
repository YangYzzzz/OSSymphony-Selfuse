"""
Initial Setup: Create a 5-slide presentation with a 6x8 table on slide 3
Task ID: impress_teach_089
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_089'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q4 2025 Regional Sales Performance"
    slide1.placeholders[1].text = "Prepared by Analytics Division — Confidential"

    # ---- Slide 2: Overview ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Executive Summary"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    body_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(4.5))
    btf = body_box.text_frame
    btf.word_wrap = True
    paragraphs_text = [
        "This report summarizes regional sales performance across all territories for Q4 2025.",
        "Overall revenue increased 12.3% compared to Q3, driven primarily by strong demand in the West and Southeast regions.",
        "Product categories showing the highest growth include Enterprise Software Licenses and Cloud Infrastructure Services.",
        "Key challenges remain in the Northeast territory where supply chain delays impacted delivery timelines.",
    ]
    for i, txt in enumerate(paragraphs_text):
        if i == 0:
            p = btf.paragraphs[0]
        else:
            p = btf.add_paragraph()
        p.text = txt
        p.space_after = Pt(12)
        for r in p.runs:
            r.font.size = Pt(16)
            r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ---- Slide 3: Data Table (the key slide) ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    # Title for the slide
    title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    ttf = title_box.text_frame
    tp = ttf.paragraphs[0]
    tp.text = "Detailed Regional Breakdown — Q4 2025"
    tp.alignment = PP_ALIGN.LEFT
    tr = tp.runs[0]
    tr.font.size = Pt(24)
    tr.font.bold = True
    tr.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # 6 columns x 8 rows table with LONG column headers (horizontal by default)
    rows, cols = 8, 6
    table_shape = slide3.shapes.add_table(
        rows, cols,
        Inches(0.5), Inches(1.3),
        Inches(12.3), Inches(5.5)
    )
    table = table_shape.table

    # Long column headers
    headers = [
        "Regional Territory Name",
        "Total Revenue Generated (USD)",
        "Year-over-Year Growth Rate (%)",
        "Customer Acquisition Count",
        "Average Transaction Value (USD)",
        "Net Promoter Score (NPS)",
    ]

    # Data rows (7 data rows)
    data = [
        ["Pacific Northwest",    "$4,523,100",  "14.2%", "312",  "$14,497", "72"],
        ["Southwest Division",   "$3,891,450",  "8.7%",  "278",  "$13,998", "68"],
        ["Great Lakes Region",   "$5,102,300",  "11.5%", "401",  "$12,724", "75"],
        ["Southeast Territory",  "$4,217,800",  "16.1%", "345",  "$12,225", "71"],
        ["Northeast Corridor",   "$3,456,200",  "3.2%",  "198",  "$17,455", "64"],
        ["Mountain West Area",   "$2,789,600",  "9.8%",  "215",  "$12,975", "69"],
        ["Central Plains Zone",  "$3,145,900",  "7.4%",  "256",  "$12,289", "66"],
    ]

    # Set column widths
    col_widths = [Inches(2.2), Inches(2.2), Inches(2.1), Inches(2.0), Inches(2.0), Inches(1.8)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Fill header row (row 0) — text is HORIZONTAL (default, no rotation)
    for c, header_text in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = header_text
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # Dark header background
        from pptx.oxml.ns import qn
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '1F3A5F'})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)

    # Fill data rows
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.text = val
            if c == 0:
                p.alignment = PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ---- Slide 4: Chart placeholder / Key Insights ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    tb4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf4 = tb4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Key Insights"
    p4.alignment = PP_ALIGN.LEFT
    r4 = p4.runs[0]
    r4.font.size = Pt(28)
    r4.font.bold = True
    r4.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    insights_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(4.5))
    itf = insights_box.text_frame
    itf.word_wrap = True
    insights = [
        "Southeast Territory showed the strongest year-over-year growth at 16.1%, driven by new enterprise contracts.",
        "Northeast Corridor underperformed with only 3.2% growth due to ongoing supply chain disruptions.",
        "Great Lakes Region generated the highest absolute revenue at $5.1M with 401 new customer acquisitions.",
        "Average transaction values were highest in the Northeast ($17,455) despite lower volume.",
        "Net Promoter Scores ranged from 64 to 75, with Great Lakes leading customer satisfaction.",
    ]
    for i, txt in enumerate(insights):
        if i == 0:
            p = itf.paragraphs[0]
        else:
            p = itf.add_paragraph()
        p.text = txt
        p.space_after = Pt(10)
        for r in p.runs:
            r.font.size = Pt(14)
            r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ---- Slide 5: Next Steps ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    tb5 = slide5.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf5 = tb5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Recommended Actions for Q1 2026"
    p5.alignment = PP_ALIGN.LEFT
    r5 = p5.runs[0]
    r5.font.size = Pt(28)
    r5.font.bold = True
    r5.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    actions_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(4.5))
    atf = actions_box.text_frame
    atf.word_wrap = True
    actions = [
        "Allocate additional sales resources to the Southeast Territory to capitalize on growth momentum.",
        "Investigate and resolve supply chain bottlenecks affecting the Northeast Corridor.",
        "Implement cross-regional best practices from Great Lakes to underperforming territories.",
        "Launch targeted customer retention programs in regions with NPS below 70.",
    ]
    for i, txt in enumerate(actions):
        if i == 0:
            p = atf.paragraphs[0]
        else:
            p = atf.add_paragraph()
        p.text = txt
        p.space_after = Pt(10)
        for r in p.runs:
            r.font.size = Pt(14)
            r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
