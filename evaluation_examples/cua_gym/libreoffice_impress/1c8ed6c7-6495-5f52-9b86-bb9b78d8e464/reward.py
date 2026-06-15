"""
Reward Script: Export presentation as PDF with each slide repeated twice
Task ID: impstruct_041
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): legacy_deck.pptx has 10 slides (doubled from original 5)
  Component 2 (0.3): Slides in pptx follow correct doubled order (A,A,B,B,C,C,...)
  Component 3 (0.2): PDF exists at /home/user/Desktop/doubled_slides.pdf with 10 pages
  Component 4 (0.2): PDF content reflects the doubled slide structure
"""

import os
import sys

WORKDIR = '/home/user'
TASK_ID = 'impstruct_041'

# Expected slide title pattern after duplication (each title appears twice consecutively)
# Original order from initial: "Quarterly Business Review", "Topic A", "Topic B", "Topic C", "Thank You"
EXPECTED_TITLES = [
    "Quarterly Business Review", "Quarterly Business Review",
    "Topic A", "Topic A",
    "Topic B", "Topic B",
    "Topic C", "Topic C",
    "Thank You", "Thank You",
]


def get_slide_title(slide):
    """Extract the first non-empty text from a slide's shapes."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                # Return only the first line (title)
                return text.split('\n')[0].strip()
    return ""


def verify_task():
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    pptx_path = os.path.join(WORKDIR, 'legacy_deck.pptx')
    pdf_path = os.path.join(WORKDIR, 'Desktop', 'doubled_slides.pdf')

    # Load pptx
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load pptx {pptx_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: PPTX has 10 slides (0.3 points)
    try:
        slide_count = len(prs.slides)
        if slide_count == 10:
            print(f"PASS: Component 1 -- PPTX has 10 slides (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 -- Expected 10 slides, found {slide_count}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Slides in doubled order (0.3 points)
    try:
        actual_titles = []
        for slide in prs.slides:
            actual_titles.append(get_slide_title(slide))

        if len(actual_titles) == 10:
            match_count = 0
            for i in range(10):
                actual = actual_titles[i]
                expected = EXPECTED_TITLES[i]
                if expected.lower() in actual.lower() or actual.lower() in expected.lower():
                    match_count += 1

            if match_count == 10:
                print(f"PASS: Component 2 -- All 10 slides in correct doubled order (0.3 pts)")
                total_score += 0.3
            elif match_count >= 8:
                partial = round(0.3 * match_count / 10, 2)
                print(f"PARTIAL: Component 2 -- {match_count}/10 slides match ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 2 -- Only {match_count}/10 slides match expected order")
                print(f"  Expected: {EXPECTED_TITLES}")
                print(f"  Actual:   {actual_titles}")
        else:
            print(f"FAIL: Component 2 -- Need 10 slides for order check, found {len(actual_titles)}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: PDF exists at correct path with 10 pages (0.2 points)
    try:
        if not os.path.exists(pdf_path):
            print(f"FAIL: Component 3 -- PDF not found at {pdf_path}")
        else:
            import fitz  # PyMuPDF
            doc = fitz.open(pdf_path)
            page_count = doc.page_count
            doc.close()

            if page_count == 10:
                print(f"PASS: Component 3 -- PDF at {pdf_path} with 10 pages (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 3 -- PDF has {page_count} pages, expected 10")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: PDF content reflects doubled slide structure (0.2 points)
    try:
        if not os.path.exists(pdf_path):
            print(f"FAIL: Component 4 -- No PDF to check content")
        else:
            import fitz  # PyMuPDF
            doc = fitz.open(pdf_path)

            # Extract text per page and check that consecutive pages share content
            # Specifically check that key titles appear in the right page pairs
            key_titles = ["Topic A", "Topic B", "Topic C", "Thank You"]
            doubled_count = 0

            # Collect all page texts
            all_text = ""
            for i in range(doc.page_count):
                all_text += doc[i].get_text()
            doc.close()

            for title in key_titles:
                count = all_text.count(title)
                if count >= 2:
                    doubled_count += 1

            if doubled_count == 4:
                print(f"PASS: Component 4 -- PDF content shows all titles doubled (0.2 pts)")
                total_score += 0.2
            elif doubled_count >= 2:
                partial = round(0.2 * doubled_count / 4, 2)
                print(f"PARTIAL: Component 4 -- {doubled_count}/4 titles doubled in PDF ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 4 -- Only {doubled_count}/4 titles appear doubled in PDF")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist app state before verification (LibreOffice may have unsaved changes)
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


persist_app_state()
verify_task()
