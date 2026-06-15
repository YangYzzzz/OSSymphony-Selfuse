"""
Initial Setup: Create a 5-slide presentation for the duplicate-and-export task
Task ID: impstruct_041
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impstruct_041'
OUTPUT = f'{WORKDIR}/legacy_deck.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    """Add a title slide (layout 0)."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points):
    """Add a content slide (layout 1) with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(18)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Intro (title slide)
    add_title_slide(
        prs,
        "Quarterly Business Review",
        "Prepared by the Strategy & Operations Team\nMarch 2025"
    )

    # Slide 2: Topic A (content slide)
    add_content_slide(prs, "Topic A", [
        "Revenue grew 18% year-over-year to $4.2M in Q1",
        "Customer acquisition cost decreased by 12% through organic channels",
        "Three new enterprise contracts signed with Fortune 500 firms",
        "Net promoter score improved from 62 to 71 across all segments",
    ])

    # Slide 3: Topic B (content slide)
    add_content_slide(prs, "Topic B", [
        "Engineering team shipped 14 major features this quarter",
        "System uptime maintained at 99.97% across all regions",
        "Migrated core infrastructure to Kubernetes with zero downtime",
        "Reduced average page load time from 2.3s to 0.8s",
    ])

    # Slide 4: Topic C (content slide)
    add_content_slide(prs, "Topic C", [
        "Marketing campaigns reached 1.2M impressions on social media",
        "Email open rates increased to 34% with new segmentation strategy",
        "Launched partner referral program with 45 active affiliates",
        "Brand awareness survey shows 28% unaided recall in target market",
    ])

    # Slide 5: End (title slide)
    add_title_slide(
        prs,
        "Thank You",
        "Questions & Discussion\nContact: strategy@acmecorp.com"
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Verify slide count
    check = Presentation(OUTPUT)
    print(f'Slide count: {len(check.slides)}')
    for i, slide in enumerate(check.slides):
        title = slide.shapes.title.text if slide.shapes.title else "(no title)"
        print(f'  Slide {i+1}: {title}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
