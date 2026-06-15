"""
Initial Setup: Chemistry Lab Report presentation with 9 slides
Task ID: impress_stu_014
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_014'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_shape(shape, text, font_size=Pt(18), bold=False, color=None, alignment=None):
    """Helper to set text on a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def create_content_slide(prs, title_text, body_lines):
    """Create a content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(18)
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Chemistry Lab Report"
    slide1.placeholders[1].text = "Acid-Base Titration of Vinegar\nPrepared by: Emily Rodriguez\nCHEM 201 - Section B\nMarch 2025"

    # --- Slide 2: Introduction ---
    create_content_slide(prs, "Introduction", [
        "Titration is a quantitative analytical technique used to determine the concentration of an unknown solution",
        "This experiment determines the acetic acid concentration in commercial white vinegar",
        "The expected concentration of acetic acid in vinegar is approximately 4-8% by mass",
        "We use NaOH as the titrant with phenolphthalein as the indicator",
    ])

    # --- Slide 3: Materials & Equipment ---
    create_content_slide(prs, "Materials & Equipment", [
        "50 mL burette (Class A, +/- 0.05 mL)",
        "250 mL Erlenmeyer flasks (3)",
        "10 mL volumetric pipette",
        "0.1000 M NaOH standardized solution",
        "Commercial white vinegar (Heinz brand)",
        "Phenolphthalein indicator solution (1% in ethanol)",
        "Distilled water, wash bottle, and ring stand assembly",
    ])

    # --- Slide 4: Experimental Procedure ---
    create_content_slide(prs, "Experimental Procedure", [
        "Step 1: Pipette 10.00 mL vinegar into a clean 250 mL Erlenmeyer flask",
        "Step 2: Add 50 mL distilled water and 3 drops phenolphthalein indicator",
        "Step 3: Fill burette with standardized 0.1000 M NaOH, record initial volume",
        "Step 4: Titrate slowly until persistent pale pink endpoint (~30 seconds)",
        "Step 5: Record final burette reading and calculate volume of NaOH used",
        "Step 6: Repeat for two additional trials",
    ])

    # --- Slide 5: Raw Data ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_box = slide5.shapes[-1]
    add_text_to_shape(title_box, "Raw Data", font_size=Pt(28), bold=True,
                      color=RGBColor(0x1F, 0x49, 0x7D))

    # Add a data table
    rows, cols = 5, 4
    tbl_shape = slide5.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(3))
    table = tbl_shape.table
    headers = ["Measurement", "Trial 1", "Trial 2", "Trial 3"]
    data = [
        ["Initial Burette (mL)", "0.15", "0.22", "0.10"],
        ["Final Burette (mL)", "17.63", "17.78", "17.55"],
        ["Volume NaOH (mL)", "17.48", "17.56", "17.45"],
        ["Vinegar Volume (mL)", "10.00", "10.00", "10.00"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(14)

    # --- Slide 6: Calculations ---
    create_content_slide(prs, "Calculations", [
        "Moles NaOH = Molarity x Volume = 0.1000 M x 0.01748 L = 0.001748 mol",
        "Mole ratio: CH3COOH : NaOH = 1 : 1",
        "Moles acetic acid = 0.001748 mol",
        "Mass acetic acid = 0.001748 mol x 60.052 g/mol = 0.1049 g",
        "Density of vinegar ~ 1.005 g/mL; Mass of 10 mL = 10.05 g",
        "Percent acetic acid = (0.1049 / 10.05) x 100 = 1.044% (per trial avg)",
    ])

    # --- Slide 7: Results Summary ---
    create_content_slide(prs, "Results Summary", [
        "Average volume NaOH used: 17.50 mL (+/- 0.06 mL)",
        "Average moles acetic acid: 0.001750 mol",
        "Calculated acetic acid concentration: 4.98% (w/w)",
        "Literature value: 5.0% acetic acid in white vinegar",
        "Percent error: 0.4%",
        "Relative standard deviation: 0.34%",
    ])

    # --- Slide 8: Discussion ---
    create_content_slide(prs, "Discussion", [
        "The experimental value of 4.98% closely matches the expected 5.0% concentration",
        "Low percent error (0.4%) indicates high accuracy in measurements",
        "Small RSD (0.34%) demonstrates excellent precision across trials",
        "Possible sources of error: parallax in burette reading, endpoint overshoot",
        "The phenolphthalein endpoint was clearly observable in all trials",
    ])

    # --- Slide 9: Conclusion ---
    create_content_slide(prs, "Conclusion", [
        "Acetic acid concentration in Heinz white vinegar determined to be 4.98% (w/w)",
        "Result is within 0.4% of the accepted literature value of 5.0%",
        "Titration technique proved reliable with high precision (RSD = 0.34%)",
        "Recommendations: use pH meter for more precise endpoint detection",
        "This method could be extended to analyze other weak acid solutions",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
