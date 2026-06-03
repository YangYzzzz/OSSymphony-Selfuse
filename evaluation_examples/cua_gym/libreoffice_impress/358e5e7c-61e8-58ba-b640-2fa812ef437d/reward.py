"""
Reward Script: Add a blank slide at the end with centered 'Thank You! Questions?' in 48pt bold
Task ID: impress_stu_014
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 10 exists (0.2 pts)
  Component 2: Slide 10 contains exact text 'Thank You! Questions?' (0.3 pts)
  Component 3: Text is 48pt bold (0.3 pts)
  Component 4: Text is center-aligned (0.2 pts)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_014'


def persist_app_state(domain: str):
    """Best-effort save via Ctrl+S in case file is open in LibreOffice."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            import time
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.text import PP_ALIGN
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    print(f"INFO: Presentation has {num_slides} slides")

    # Component 1: Slide 10 exists (0.2 points)
    # Initial has 9 slides; golden has 10. This checks the task-introduced change.
    try:
        if num_slides >= 10:
            print(f"PASS: Component 1 — Slide 10 exists (total slides: {num_slides}) (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 1 — Expected >= 10 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Early exit if slide 10 doesn't exist
    if num_slides < 10:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    slide10 = prs.slides[9]  # 0-indexed

    # Gather all text shapes on slide 10 (including inside groups)
    def get_all_text_shapes(slide):
        def extract(shape):
            results = []
            if hasattr(shape, "text") and hasattr(shape, "text_frame"):
                results.append(shape)
            if hasattr(shape, 'shapes'):
                for sub in shape.shapes:
                    results.extend(extract(sub))
            return results
        out = []
        for shape in slide.shapes:
            out.extend(extract(shape))
        return out

    text_shapes = get_all_text_shapes(slide10)

    # Find the shape containing the target text
    target_text = "Thank You! Questions?"
    found_shape = None
    for shape in text_shapes:
        full_text = shape.text_frame.text.strip()
        if target_text in full_text:
            found_shape = shape
            break

    # Component 2: Slide 10 contains exact text 'Thank You! Questions?' (0.3 points)
    try:
        if found_shape is not None:
            actual_text = found_shape.text_frame.text.strip()
            if actual_text == target_text:
                print(f"PASS: Component 2 — Found exact text '{target_text}' on slide 10 (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 2 — Text contains target but has extra content: {repr(actual_text)}")
        else:
            all_texts = [s.text_frame.text.strip() for s in text_shapes if s.text_frame.text.strip()]
            print(f"FAIL: Component 2 — '{target_text}' not found on slide 10. Found texts: {all_texts}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Components 3 and 4 require the found shape
    if found_shape is None:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 3: Text is 48pt bold (0.3 points)
    # Check all non-empty runs in the shape
    try:
        all_runs = []
        for para in found_shape.text_frame.paragraphs:
            for run in para.runs:
                if (run.text or "").strip():
                    all_runs.append(run)

        if not all_runs:
            print("FAIL: Component 3 — No non-empty runs found in the target text shape")
        else:
            all_bold = True
            all_48pt = True
            for run in all_runs:
                # Check bold: None means inherit (not explicitly bold)
                if run.font.bold is not True:
                    all_bold = False
                # Check font size: 48pt = 609600 EMU
                if run.font.size is None or abs(run.font.size - 609600) > 100:
                    all_48pt = False

            if all_bold and all_48pt:
                print(f"PASS: Component 3 — All runs are 48pt bold (0.3 pts)")
                total_score += 0.3
            else:
                sizes = [run.font.size for run in all_runs]
                bolds = [run.font.bold for run in all_runs]
                print(f"FAIL: Component 3 — bold={bolds}, sizes={sizes} (expected all True, all 609600)")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Text is center-aligned (0.2 points)
    try:
        all_centered = True
        para_count = 0
        for para in found_shape.text_frame.paragraphs:
            # Only check paragraphs with actual text
            nonempty = [r for r in para.runs if (r.text or "").strip()]
            if nonempty:
                para_count += 1
                if para.alignment != PP_ALIGN.CENTER:
                    all_centered = False

        if para_count == 0:
            print("FAIL: Component 4 — No non-empty paragraphs found")
        elif all_centered:
            print(f"PASS: Component 4 — Text is center-aligned (0.2 pts)")
            total_score += 0.2
        else:
            alignments = []
            for para in found_shape.text_frame.paragraphs:
                nonempty = [r for r in para.runs if (r.text or "").strip()]
                if nonempty:
                    alignments.append(para.alignment)
            print(f"FAIL: Component 4 — Alignment={alignments}, expected CENTER (2)")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
