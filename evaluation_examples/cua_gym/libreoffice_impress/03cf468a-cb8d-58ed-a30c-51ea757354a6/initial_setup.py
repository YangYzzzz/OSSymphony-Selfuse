"""
Initial Setup: Brand refresh presentation with old company name and brand colors
Task ID: impress_gf5_029
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import copy

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_029'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

OLD_BRAND_BLUE = RGBColor(0x1E, 0x90, 0xFF)  # #1E90FF
OLD_COMPANY = "TechCorp Inc."

def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                alignment=None, font_color=None, font_name="Arial"):
    """Helper to add a textbox with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if font_color:
        run.font.color.rgb = font_color
    return txBox


def add_filled_shape(slide, left, top, width, height, fill_color):
    """Add a rectangle shape with a solid fill color."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_group_with_text(slide, prs, left, top, text, fill_color):
    """Add a textbox with fill to simulate grouped content.
    We create a shape with text and fill to represent branded grouped elements."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.5), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return shape


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ==================== SLIDE DATA ====================
    slide_data = [
        # Slide 1: Title Slide
        {
            "type": "title",
            "title": f"{OLD_COMPANY} Annual Strategy Review",
            "subtitle": "Fiscal Year 2025-2026 Planning Session\nConfidential - Internal Use Only",
            "has_brand_shapes": True,
        },
        # Slide 2: Agenda
        {
            "type": "content",
            "title": f"{OLD_COMPANY} Meeting Agenda",
            "body": "1. Company Performance Overview\n2. Market Analysis & Competitive Landscape\n3. Product Roadmap Updates\n4. Financial Projections\n5. Team Expansion Plans\n6. Q&A Session",
            "has_brand_shapes": True,
        },
        # Slide 3: Company Overview
        {
            "type": "content",
            "title": "About " + OLD_COMPANY,
            "body": f"{OLD_COMPANY} is a leading enterprise software provider specializing in cloud-native solutions for Fortune 500 companies. Founded in 2012, we have grown to serve over 2,400 clients globally.",
            "has_brand_shapes": False,
        },
        # Slide 4: Revenue Dashboard
        {
            "type": "content",
            "title": f"{OLD_COMPANY} Revenue Performance",
            "body": "Q1 2025: $42.3M (+18% YoY)\nQ2 2025: $47.1M (+22% YoY)\nQ3 2025: $51.8M (+25% YoY)\nQ4 2025: $58.2M (+31% YoY)\nTotal FY2025: $199.4M",
            "has_brand_shapes": True,
        },
        # Slide 5: Market Position
        {
            "type": "content",
            "title": "Market Position",
            "body": f"{OLD_COMPANY} holds a 14.7% market share in the enterprise middleware segment. Key competitors include DataBridge Solutions (11.2%) and CloudPath Technologies (9.8%).",
            "has_brand_shapes": False,
        },
        # Slide 6: Product Portfolio
        {
            "type": "content",
            "title": f"{OLD_COMPANY} Product Suite",
            "body": "CorePlatform - Enterprise Integration Hub\nDataSync Pro - Real-time Data Pipeline\nSecureGate - Zero Trust Security Layer\nAnalytix - Business Intelligence Dashboard\nDevOps Forge - CI/CD Automation",
            "has_brand_shapes": True,
        },
        # Slide 7: Customer Success
        {
            "type": "content",
            "title": "Customer Success Stories",
            "body": f"GlobalBank Corp reduced processing time by 67% using {OLD_COMPANY} CorePlatform.\nRetailMax achieved 99.99% uptime with {OLD_COMPANY} SecureGate.\nMediHealth improved data accuracy by 43% with DataSync Pro.",
            "has_brand_shapes": False,
        },
        # Slide 8: Team Growth
        {
            "type": "content",
            "title": f"{OLD_COMPANY} Team Expansion",
            "body": "Current Headcount: 1,247 employees\nEngineering: 612 (+85 planned)\nSales & Marketing: 298 (+42 planned)\nCustomer Success: 187 (+31 planned)\nOperations: 150 (+18 planned)",
            "has_brand_shapes": True,
        },
        # Slide 9: Technology Stack
        {
            "type": "content",
            "title": "Technology Infrastructure",
            "body": f"All {OLD_COMPANY} products run on a microservices architecture deployed across AWS, Azure, and GCP. Our proprietary container orchestration layer handles 2.3 billion API calls daily.",
            "has_brand_shapes": True,
        },
        # Slide 10: R&D Investments
        {
            "type": "content",
            "title": f"{OLD_COMPANY} R&D Focus Areas",
            "body": "AI/ML Integration: $28M investment\nEdge Computing: $15M investment\nQuantum-Ready Security: $8M investment\nGreen Computing Initiative: $5M investment",
            "has_brand_shapes": False,
        },
        # Slide 11: Partnership Ecosystem
        {
            "type": "content",
            "title": "Strategic Partnerships",
            "body": f"{OLD_COMPANY} maintains tier-1 partnerships with:\n- Amazon Web Services (Advanced Partner)\n- Microsoft Azure (Gold Partner)\n- Google Cloud (Premier Partner)\n- Salesforce (ISV Partner)",
            "has_brand_shapes": True,
        },
        # Slide 12: Global Presence
        {
            "type": "content",
            "title": f"{OLD_COMPANY} Global Offices",
            "body": "San Francisco, CA (HQ)\nNew York, NY\nLondon, UK\nSingapore\nTokyo, Japan\nSydney, Australia\nBerlin, Germany",
            "has_brand_shapes": True,
        },
        # Slide 13: Financial Outlook
        {
            "type": "content",
            "title": "FY2026 Financial Projections",
            "body": f"{OLD_COMPANY} projects revenue of $265M for FY2026, representing 33% growth. EBITDA margin expected to improve from 22% to 27%. Free cash flow projected at $58M.",
            "has_brand_shapes": False,
        },
        # Slide 14: Product Roadmap
        {
            "type": "content",
            "title": f"{OLD_COMPANY} Product Roadmap 2026",
            "body": "Q1: CorePlatform 4.0 with AI Assistant\nQ2: DataSync Pro Edge Edition\nQ3: SecureGate Quantum Module\nQ4: Unified Analytics Dashboard v2",
            "has_brand_shapes": True,
        },
        # Slide 15: Customer Metrics
        {
            "type": "content",
            "title": "Customer Satisfaction Metrics",
            "body": f"{OLD_COMPANY} Net Promoter Score: 72 (Industry avg: 45)\nCustomer Retention Rate: 96.3%\nAverage Contract Value: $83,500\nTime to Value: 14 days (down from 28)",
            "has_brand_shapes": False,
        },
        # Slide 16: Security & Compliance
        {
            "type": "content",
            "title": f"{OLD_COMPANY} Security Standards",
            "body": "SOC 2 Type II Certified\nISO 27001 Compliant\nGDPR & CCPA Ready\nFedRAMP Authorized\nHIPAA Compliant\nPCI DSS Level 1",
            "has_brand_shapes": True,
        },
        # Slide 17: Sustainability
        {
            "type": "content",
            "title": "Sustainability Initiatives",
            "body": f"{OLD_COMPANY} committed to carbon neutrality by 2027.\nData center PUE reduced to 1.15.\nRemote-first policy reduces commute emissions by 62%.\n100% renewable energy in all offices.",
            "has_brand_shapes": True,
        },
        # Slide 18: Awards & Recognition
        {
            "type": "content",
            "title": f"{OLD_COMPANY} Awards 2025",
            "body": "Gartner Magic Quadrant Leader - Enterprise Integration\nForbes Cloud 100 - Ranked #34\nBest Places to Work - Glassdoor 4.6/5\nFast Company Most Innovative - #12",
            "has_brand_shapes": False,
        },
        # Slide 19: Action Items
        {
            "type": "content",
            "title": "Key Action Items",
            "body": f"1. Finalize {OLD_COMPANY} brand refresh by March 15\n2. Launch CorePlatform 4.0 beta program\n3. Complete Series D funding round\n4. Expand Singapore engineering hub\n5. Roll out new partner certification program",
            "has_brand_shapes": True,
        },
        # Slide 20: Closing
        {
            "type": "title",
            "title": f"Thank You\n{OLD_COMPANY}",
            "subtitle": "Questions? Contact: strategy@techcorp.com\nwww.techcorp-inc.com",
            "has_brand_shapes": True,
        },
    ]

    for i, sd in enumerate(slide_data):
        if sd["type"] == "title":
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank

            # Title text
            add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(2),
                        sd["title"], font_size=36, bold=True,
                        alignment=PP_ALIGN.CENTER, font_color=RGBColor(0x33, 0x33, 0x33))

            # Subtitle
            add_textbox(slide, Inches(2), Inches(4), Inches(9), Inches(1.5),
                        sd["subtitle"], font_size=18,
                        alignment=PP_ALIGN.CENTER, font_color=RGBColor(0x66, 0x66, 0x66))

            # Brand color bar at top
            add_filled_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.4), OLD_BRAND_BLUE)

            # Brand color bar at bottom
            add_filled_shape(slide, Inches(0), Inches(7.1), prs.slide_width, Inches(0.4), OLD_BRAND_BLUE)

            if sd.get("has_brand_shapes"):
                # Brand accent shape
                add_filled_shape(slide, Inches(5.5), Inches(6.0), Inches(2.5), Inches(0.08), OLD_BRAND_BLUE)

        else:  # content slide
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank

            # Header bar with brand color
            add_filled_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.0), OLD_BRAND_BLUE)

            # Title on header bar
            add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                        sd["title"], font_size=28, bold=True,
                        alignment=PP_ALIGN.LEFT, font_color=RGBColor(0xFF, 0xFF, 0xFF))

            # Body content
            add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(5),
                        sd["body"], font_size=16,
                        alignment=PP_ALIGN.LEFT, font_color=RGBColor(0x33, 0x33, 0x33))

            # Footer with company name
            add_textbox(slide, Inches(0.5), Inches(6.8), Inches(5), Inches(0.5),
                        f"Confidential - {OLD_COMPANY}", font_size=10,
                        alignment=PP_ALIGN.LEFT, font_color=RGBColor(0x99, 0x99, 0x99))

            if sd.get("has_brand_shapes"):
                # Brand accent sidebar
                add_filled_shape(slide, Inches(12.5), Inches(1.2), Inches(0.15), Inches(5.5), OLD_BRAND_BLUE)

                # Branded badge shape with company name
                add_group_with_text(slide, prs, Inches(9.5), Inches(5.8),
                                    OLD_COMPANY, OLD_BRAND_BLUE)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # Count items for reference
    text_count = 0
    color_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text') and OLD_COMPANY in (shape.text or ''):
                text_count += 1
            if shape.fill and shape.fill.type is not None:
                try:
                    if shape.fill.fore_color and shape.fill.fore_color.rgb == OLD_BRAND_BLUE:
                        color_count += 1
                except:
                    pass
    print(f'Shapes with "{OLD_COMPANY}": {text_count}')
    print(f'Shapes with old brand blue: {color_count}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
