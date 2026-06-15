"""
Reward Script: Brand refresh batch update across 20-slide presentation
Task ID: impress_gf5_029
Domain: libreoffice_impress
Scoring:
  Component 1 (0.35): All 'TechCorp Inc.' replaced with 'Innovate Systems LLC'
  Component 2 (0.35): All old blue (#1E90FF) fills replaced with new blue (#0052CC)
  Component 3 (0.15): No residual 'TechCorp Inc.' text remains anywhere
  Component 4 (0.15): No residual old blue (#1E90FF) fills remain anywhere
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_029'


def persist_app_state(domain: str):
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_all_shapes_recursive(shapes):
    """Recursively get all shapes including those inside groups."""
    results = []
    for shape in shapes:
        results.append(shape)
        if hasattr(shape, 'shapes'):
            results.extend(get_all_shapes_recursive(shape.shapes))
    return results


def verify_task(file_path):
    """
    Verify brand refresh task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have 20 slides
    if len(prs.slides) != 20:
        print(f"PRECONDITION FAIL: Expected 20 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Gather counts across all slides
    techcorp_count = 0
    innovate_count = 0
    old_blue_count = 0
    new_blue_count = 0

    for slide in prs.slides:
        all_shapes = get_all_shapes_recursive(slide.shapes)
        for shape in all_shapes:
            # Check text
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text
                    if 'TechCorp Inc.' in text:
                        techcorp_count += 1
                    if 'Innovate Systems LLC' in text:
                        innovate_count += 1

            # Check fill color
            if hasattr(shape, 'fill'):
                try:
                    fill = shape.fill
                    if fill.type is not None and fill.type == 1:  # solid fill
                        rgb = str(fill.fore_color.rgb).upper()
                        if rgb == '1E90FF':
                            old_blue_count += 1
                        elif rgb == '0052CC':
                            new_blue_count += 1
                except Exception:
                    pass

    print(f"INFO: TechCorp Inc. mentions: {techcorp_count}")
    print(f"INFO: Innovate Systems LLC mentions: {innovate_count}")
    print(f"INFO: Old blue (#1E90FF) fills: {old_blue_count}")
    print(f"INFO: New blue (#0052CC) fills: {new_blue_count}")

    # Component 1: 'Innovate Systems LLC' text is present (0.35 points)
    # This checks that the replacement text exists — should FAIL on initial (0 mentions)
    # and PASS on golden (50 mentions).
    try:
        if innovate_count > 0:
            print(f"PASS: Component 1 — Found {innovate_count} 'Innovate Systems LLC' mentions (0.35 pts)")
            total_score += 0.35
        else:
            print(f"FAIL: Component 1 — No 'Innovate Systems LLC' text found")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: New brand blue (#0052CC) fills are present (0.35 points)
    # Should FAIL on initial (0 new blue fills) and PASS on golden (46 fills).
    try:
        if new_blue_count > 0:
            print(f"PASS: Component 2 — Found {new_blue_count} new blue (#0052CC) fills (0.35 pts)")
            total_score += 0.35
        else:
            print(f"FAIL: Component 2 — No new blue (#0052CC) fills found")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: No residual 'TechCorp Inc.' text (0.15 points)
    # Should FAIL on initial (50 mentions) and PASS on golden (0 mentions).
    try:
        if techcorp_count == 0:
            print(f"PASS: Component 3 — No residual 'TechCorp Inc.' text found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 3 — Found {techcorp_count} residual 'TechCorp Inc.' mentions")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: No residual old blue (#1E90FF) fills (0.15 points)
    # Should FAIL on initial (46 old blue fills) and PASS on golden (0 old blue fills).
    try:
        if old_blue_count == 0:
            print(f"PASS: Component 4 — No residual old blue (#1E90FF) fills (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 — Found {old_blue_count} residual old blue (#1E90FF) fills")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
