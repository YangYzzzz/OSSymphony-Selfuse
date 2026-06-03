"""
Reward Script: Import master slide from Template_Source.pptx into My_Slides.pptx
Task ID: impress_ma_017
Domain: libreoffice_impress
Scoring:
  Component 1: A slide master with gradient background exists (0.30 pts)
  Component 2: That master has Montserrat fonts in placeholders (0.30 pts)
  Component 3: That master has a logo/picture shape (0.20 pts)
  Component 4: Original slides still intact - 6 slides present (0.20 pts)
             (anchored to the imported master existing via compound check)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_017'


def persist_app_state(domain):
    """Save any unsaved GUI edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify that My_Slides.pptx now contains the master slide imported from
    Template_Source.pptx, with gradient background, Montserrat fonts, and logo.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: file must be loadable
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Find which slide masters have key properties by scanning ZIP XML
    # We need to find a master with gradient background (gradFill)
    gradient_master = None
    montserrat_in_gradient_master = False
    pic_in_gradient_master = False

    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            master_files = sorted([f for f in zf.namelist()
                                   if 'slideMasters/slideMaster' in f
                                   and f.endswith('.xml')
                                   and '_rels' not in f])
            print(f"INFO: Master XML files: {master_files}")

            for mf in master_files:
                with zf.open(mf) as f:
                    content = f.read().decode('utf-8')
                    has_gradient = 'gradFill' in content
                    has_montserrat = 'Montserrat' in content
                    has_pic = '<p:pic' in content
                    print(f"INFO: {mf}: gradient={has_gradient}, montserrat={has_montserrat}, pic={has_pic}")

                    if has_gradient:
                        gradient_master = mf
                        montserrat_in_gradient_master = has_montserrat
                        pic_in_gradient_master = has_pic
    except Exception as e:
        print(f"ERROR: Cannot scan ZIP contents: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: A slide master with gradient background exists (0.30 pts)
    # This is the primary indicator that the Template_Source master was imported.
    # Initial state: no master has gradient. Golden: imported master has gradient.
    try:
        if gradient_master is not None:
            print(f"PASS: Component 1 - Slide master with gradient background found in {gradient_master} (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 1 - No slide master with gradient background found")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: The gradient master has Montserrat fonts (0.30 pts)
    # The Template_Source master uses Montserrat. This must be preserved on import.
    try:
        if gradient_master is not None and montserrat_in_gradient_master:
            print(f"PASS: Component 2 - Gradient master has Montserrat fonts (0.30 pts)")
            total_score += 0.30
        else:
            if gradient_master is None:
                print(f"FAIL: Component 2 - No gradient master found, cannot check fonts")
            else:
                print(f"FAIL: Component 2 - Gradient master does not have Montserrat fonts")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: The gradient master has a logo/picture shape (0.20 pts)
    # Template_Source has a Logo picture shape on its master.
    try:
        if gradient_master is not None and pic_in_gradient_master:
            print(f"PASS: Component 3 - Gradient master has logo/picture shape (0.20 pts)")
            total_score += 0.20
        else:
            if gradient_master is None:
                print(f"FAIL: Component 3 - No gradient master found, cannot check logo")
            else:
                print(f"FAIL: Component 3 - Gradient master does not have picture/logo shape")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Original slides still intact AND imported master exists (0.20 pts)
    # Compound check: presentation has 6 slides AND the gradient master exists.
    # This ensures the import didn't break the existing content.
    try:
        slide_count = len(prs.slides)
        if gradient_master is not None and slide_count == 6:
            print(f"PASS: Component 4 - 6 slides intact with imported master present (0.20 pts)")
            total_score += 0.20
        else:
            if gradient_master is None:
                print(f"FAIL: Component 4 - No imported master, compound check fails")
            else:
                print(f"FAIL: Component 4 - Expected 6 slides, found {slide_count}")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
