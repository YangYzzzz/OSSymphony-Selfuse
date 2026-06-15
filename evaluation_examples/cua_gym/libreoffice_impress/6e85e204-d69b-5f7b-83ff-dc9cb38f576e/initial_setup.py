"""
Initial Setup: Copy master slide from Template_Source.pptx into My_Slides.pptx
Task ID: impress_ma_017
Domain: libreoffice_impress

Creates two files:
  - /home/user/Template_Source.pptx  (professionally designed master with blue gradient, Montserrat, logo)
  - /home/user/impress_ma_017.pptx   (6 slides with a plain white master)
Opens impress_ma_017.pptx in LibreOffice Impress.
"""

import os
import shlex
import subprocess
import time
import zipfile
import shutil
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_017'
MY_SLIDES = f'{WORKDIR}/{TASK_ID}.pptx'
TEMPLATE_SOURCE = f'{WORKDIR}/Template_Source.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_logo_image(path):
    """Create a simple logo PNG image."""
    from PIL import Image, ImageDraw, ImageFont
    img = Image.new('RGBA', (200, 80), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    # Blue circle
    draw.ellipse([10, 10, 60, 60], fill=(30, 80, 180, 255))
    # White letter inside circle
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 28)
    except:
        font = ImageFont.load_default()
    draw.text((24, 16), "A", fill=(255, 255, 255, 255), font=font)
    # Company name next to circle
    try:
        font2 = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 22)
    except:
        font2 = ImageFont.load_default()
    draw.text((70, 22), "Acuity Corp", fill=(30, 80, 180, 255), font=font2)
    img.save(path)
    return path


def create_template_source():
    """Create Template_Source.pptx with a professionally designed master slide."""
    prs = Presentation()

    # Access the slide master and modify it
    master = prs.slide_masters[0]

    # Set gradient background on master via XML manipulation
    bg = master.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = RGBColor(0x0D, 0x47, 0xA1)  # dark blue
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = RGBColor(0x42, 0xA5, 0xF5)  # lighter blue
    fill.gradient_stops[1].position = 1.0

    # Modify master placeholders to use Montserrat font
    for shape in master.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = "Montserrat"
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Also set font on all slide layouts
    for layout in prs.slide_layouts:
        for shape in layout.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = "Montserrat"
                        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    # Set default font on paragraph level via XML
                    pPr = para._p.get_or_add_pPr()
                    from pptx.oxml.ns import qn
                    defRPr = pPr.find(qn('a:defRPr'))
                    if defRPr is None:
                        defRPr = pPr.makeelement(qn('a:defRPr'), {})
                        pPr.append(defRPr)
                    # Set latin font
                    latin = defRPr.find(qn('a:latin'))
                    if latin is None:
                        latin = defRPr.makeelement(qn('a:latin'), {'typeface': 'Montserrat'})
                        defRPr.append(latin)
                    else:
                        latin.set('typeface', 'Montserrat')

    # Add a title text to first layout slide as a demo
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Template Presentation"
    slide.placeholders[1].text = "Professional Design Template"
    for para in slide.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Montserrat"
            run.font.size = Pt(36)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    for para in slide.placeholders[1].text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Montserrat"
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

    # Also add logo directly to the slide (since master/layout shapes can't add_picture)
    logo_path = f'{WORKDIR}/_temp_logo.png'
    create_logo_image(logo_path)
    slide.shapes.add_picture(logo_path, Inches(0.3), Inches(0.2), Inches(1.8), Inches(0.6))

    prs.save(TEMPLATE_SOURCE)

    # Now inject the logo into the slide master via ZIP/XML manipulation
    _inject_logo_into_master(TEMPLATE_SOURCE, logo_path)

    # Clean up temp logo
    if os.path.exists(logo_path):
        os.remove(logo_path)

    print(f'Template source created: {TEMPLATE_SOURCE}')


def _inject_logo_into_master(pptx_path, logo_path):
    """Inject a logo image into the slide master via ZIP-level XML manipulation."""
    import zipfile
    import xml.etree.ElementTree as ET

    tmp_path = pptx_path + '.tmp'

    with open(logo_path, 'rb') as f:
        logo_bytes = f.read()

    with zipfile.ZipFile(pptx_path, 'r') as zin:
        files = {}
        for name in zin.namelist():
            files[name] = zin.read(name)

    # Find next available image number
    media_files = [n for n in files if n.startswith('ppt/media/')]
    next_img = len(media_files) + 1
    new_media_name = f'image{next_img}.png'
    new_media_path = f'ppt/media/{new_media_name}'
    files[new_media_path] = logo_bytes

    # Add relationship in slideMaster1.xml.rels
    master_rels_path = 'ppt/slideMasters/_rels/slideMaster1.xml.rels'
    rels_xml = files[master_rels_path].decode('utf-8')
    rels_root = ET.fromstring(rels_xml)

    # Find max rId
    max_rid = 0
    for rel in rels_root:
        rid = rel.get('Id', '')
        if rid.startswith('rId'):
            try:
                num = int(rid[3:])
                if num > max_rid:
                    max_rid = num
            except ValueError:
                pass
    new_rid = f'rId{max_rid + 1}'

    # Add image relationship
    new_rel = ET.SubElement(rels_root, 'Relationship')
    new_rel.set('Id', new_rid)
    new_rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image')
    new_rel.set('Target', f'../media/{new_media_name}')

    files[master_rels_path] = ET.tostring(rels_root, xml_declaration=True, encoding='UTF-8')

    # Add picture shape to slideMaster1.xml
    master_xml = files['ppt/slideMasters/slideMaster1.xml'].decode('utf-8')
    master_root = ET.fromstring(master_xml)

    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

    ET.register_namespace('', ns_p)
    ET.register_namespace('a', ns_a)
    ET.register_namespace('r', ns_r)
    ET.register_namespace('p', ns_p)

    # Find cSld/spTree
    cSld = master_root.find(f'{{{ns_p}}}cSld')
    spTree = cSld.find(f'{{{ns_p}}}spTree') if cSld is not None else None

    if spTree is not None:
        # Create a pic element for the logo
        # EMU values: left=0.3in=274320, top=0.2in=182880, width=1.8in=1645920, height=0.6in=548640
        pic_xml = f'''<p:pic xmlns:p="{ns_p}" xmlns:a="{ns_a}" xmlns:r="{ns_r}">
          <p:nvPicPr>
            <p:cNvPr id="9999" name="Logo"/>
            <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
            <p:nvPr/>
          </p:nvPicPr>
          <p:blipFill>
            <a:blip r:embed="{new_rid}"/>
            <a:stretch><a:fillRect/></a:stretch>
          </p:blipFill>
          <p:spPr>
            <a:xfrm>
              <a:off x="274320" y="182880"/>
              <a:ext cx="1645920" cy="548640"/>
            </a:xfrm>
            <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
          </p:spPr>
        </p:pic>'''
        pic_elem = ET.fromstring(pic_xml)
        spTree.append(pic_elem)

    files['ppt/slideMasters/slideMaster1.xml'] = ET.tostring(master_root, xml_declaration=True, encoding='UTF-8')

    # Update Content_Types if needed
    ct_xml = files['[Content_Types].xml'].decode('utf-8')
    if '.png' not in ct_xml.lower() or 'image/png' not in ct_xml:
        ct_root = ET.fromstring(ct_xml)
        ct_ns = 'http://schemas.openxmlformats.org/package/2006/content-types'
        ET.register_namespace('', ct_ns)
        # Check if png extension already registered
        has_png = False
        for elem in ct_root:
            if elem.get('Extension', '').lower() == 'png':
                has_png = True
                break
        if not has_png:
            ext = ET.SubElement(ct_root, f'{{{ct_ns}}}Default')
            ext.set('Extension', 'png')
            ext.set('ContentType', 'image/png')
        files['[Content_Types].xml'] = ET.tostring(ct_root, xml_declaration=True, encoding='UTF-8')

    # Write new ZIP
    with zipfile.ZipFile(tmp_path, 'w', zipfile.ZIP_DEFLATED) as zout:
        for name, data in files.items():
            if isinstance(data, str):
                data = data.encode('utf-8')
            zout.writestr(name, data)

    shutil.move(tmp_path, pptx_path)


def create_my_slides():
    """Create My_Slides (impress_ma_017.pptx) with 6 slides and a plain white master."""
    prs = Presentation()

    # Slide 1: Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q4 2025 Business Review"
    slide1.placeholders[1].text = "Acuity Corporation - Internal Report"
    for para in slide1.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(36)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    for para in slide1.placeholders[1].text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # Slide 2: Revenue Overview
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(1))
    tf2 = title2.text_frame
    p = tf2.paragraphs[0]
    p.text = "Revenue Overview"
    p.runs[0].font.name = "Calibri"
    p.runs[0].font.size = Pt(28)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    table2 = slide2.shapes.add_table(5, 4, Inches(0.8), Inches(1.8), Inches(8), Inches(3)).table
    headers = ["Region", "Q3 Revenue", "Q4 Revenue", "Growth"]
    data = [
        ["North America", "$12.4M", "$14.1M", "+13.7%"],
        ["Europe", "$8.7M", "$9.2M", "+5.7%"],
        ["Asia Pacific", "$6.3M", "$7.8M", "+23.8%"],
        ["Latin America", "$3.1M", "$3.5M", "+12.9%"],
    ]
    for c, h in enumerate(headers):
        table2.cell(0, c).text = h
    for r, row in enumerate(data, 1):
        for c, val in enumerate(row):
            table2.cell(r, c).text = val

    # Slide 3: Product Performance
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    title3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(1))
    tf3 = title3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Product Performance"
    p3.runs[0].font.name = "Calibri"
    p3.runs[0].font.size = Pt(28)
    p3.runs[0].font.bold = True
    p3.runs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    body3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8), Inches(4))
    tf_body3 = body3.text_frame
    tf_body3.word_wrap = True
    products = [
        ("Cloud Platform", "Revenue grew 28% YoY. Enterprise adoption surged with 145 new accounts."),
        ("Analytics Suite", "Maintained steady growth at 12% YoY. New ML features drove upgrades."),
        ("Security Module", "Fastest growing product at 42% YoY. SOC2 compliance features popular."),
    ]
    for i, (name, desc) in enumerate(products):
        if i == 0:
            p = tf_body3.paragraphs[0]
        else:
            p = tf_body3.add_paragraph()
        run = p.add_run() if i > 0 or len(p.runs) == 0 else p.runs[0]
        run.text = name
        run.font.name = "Calibri"
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p2 = tf_body3.add_paragraph()
        run2 = p2.add_run()
        run2.text = desc
        run2.font.name = "Calibri"
        run2.font.size = Pt(14)
        run2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # Slide 4: Team Highlights
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    title4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(1))
    tf4 = title4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Team Highlights"
    p4.runs[0].font.name = "Calibri"
    p4.runs[0].font.size = Pt(28)
    p4.runs[0].font.bold = True
    p4.runs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    body4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8), Inches(4))
    tf_body4 = body4.text_frame
    tf_body4.word_wrap = True
    highlights = [
        "Engineering headcount grew to 342 (up from 298 in Q3)",
        "Customer Success team achieved 96.2% retention rate",
        "Marketing launched the 'Acuity Insights' webinar series (12K registrations)",
        "Sales closed 3 enterprise deals over $2M each",
    ]
    for i, text in enumerate(highlights):
        if i == 0:
            p = tf_body4.paragraphs[0]
        else:
            p = tf_body4.add_paragraph()
        p.text = text
        p.level = 0
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # Slide 5: Key Initiatives for Q1 2026
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    title5 = slide5.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(1))
    tf5 = title5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Key Initiatives for Q1 2026"
    p5.runs[0].font.name = "Calibri"
    p5.runs[0].font.size = Pt(28)
    p5.runs[0].font.bold = True
    p5.runs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    table5 = slide5.shapes.add_table(5, 3, Inches(0.8), Inches(1.8), Inches(8), Inches(3)).table
    headers5 = ["Initiative", "Owner", "Target Date"]
    data5 = [
        ["Platform v3.0 Launch", "Sarah Chen", "Feb 15, 2026"],
        ["APAC Market Expansion", "Marcus Johnson", "Mar 1, 2026"],
        ["SOC2 Type II Certification", "Li Wei", "Jan 31, 2026"],
        ["Partner Program Rollout", "Priya Sharma", "Mar 15, 2026"],
    ]
    for c, h in enumerate(headers5):
        table5.cell(0, c).text = h
    for r, row in enumerate(data5, 1):
        for c, val in enumerate(row):
            table5.cell(r, c).text = val

    # Slide 6: Thank You
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    thanks = slide6.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(2))
    tf6 = thanks.text_frame
    p6 = tf6.paragraphs[0]
    p6.text = "Thank You"
    p6.alignment = PP_ALIGN.CENTER
    p6.runs[0].font.name = "Calibri"
    p6.runs[0].font.size = Pt(40)
    p6.runs[0].font.bold = True
    p6.runs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p_sub = tf6.add_paragraph()
    p_sub.text = "Questions? Contact leadership@acuity.com"
    p_sub.alignment = PP_ALIGN.CENTER
    run_sub = p_sub.runs[0]
    run_sub.font.name = "Calibri"
    run_sub.font.size = Pt(16)
    run_sub.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    prs.save(MY_SLIDES)
    print(f'My_Slides created: {MY_SLIDES}')


def main():
    create_template_source()
    create_my_slides()

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{MY_SLIDES}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


main()
