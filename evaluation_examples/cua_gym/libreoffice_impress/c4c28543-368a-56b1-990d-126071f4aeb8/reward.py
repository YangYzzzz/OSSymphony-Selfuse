"""
Reward Script: Extract presenter notes from Strategy_2025.pptx into strategy_notes.docx
Task ID: osworld_multi_apps_impress_notes_export_009
Domain: libreoffice_impress + libreoffice_writer (multi-app)

Scoring Rubric:
  Precondition: strategy_notes.docx exists on Desktop (gate — returns 0.0 if absent)
  Component 1: All 14 'Slide N:' labels present in correct order (0.4 pts)
  Component 2: Notes text content matches pptx notes from Strategy_2025.pptx (0.4 pts)
  Component 3: Notes sections separated by blank paragraph lines (0.2 pts)
  Total: 1.0
"""

import os
import re

# python-pptx for reading notes from pptx
from pptx import Presentation
# python-docx for reading the output docx
from docx import Document

WORKDIR = '/home/user'
TASK_ID = 'osworld_multi_apps_impress_notes_export_009'

PPTX_PATH = f'{WORKDIR}/Desktop/Strategy_2025.pptx'
DOCX_PATH = f'{WORKDIR}/Desktop/strategy_notes.docx'


def get_slide_notes(pptx_path):
    """Extract all notes from pptx slides, returns list of (slide_num, notes_text)."""
    prs = Presentation(pptx_path)
    notes = []
    for i, slide in enumerate(prs.slides):
        try:
            text = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            text = ""
        notes.append((i + 1, text))
    return notes


def verify_task():
    """
    Verify that strategy_notes.docx was created with correct structure.
    Returns float between 0.0 and 1.0.
    """
    total_score = 0.0

    # Precondition gate: strategy_notes.docx must exist on Desktop
    if not os.path.exists(DOCX_PATH):
        print(f"FAIL: strategy_notes.docx not found at {DOCX_PATH}")
        print("REWARD: 0.0")
        return 0.0

    # Load the docx
    try:
        doc = Document(DOCX_PATH)
    except Exception as e:
        print(f"CRITICAL: Cannot load {DOCX_PATH}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Load source pptx notes (for content verification)
    try:
        slide_notes = get_slide_notes(PPTX_PATH)
        num_slides = len(slide_notes)
    except Exception as e:
        print(f"CRITICAL: Cannot load {PPTX_PATH}: {e}")
        print("REWARD: 0.0")
        return 0.0

    paragraphs = [p.text for p in doc.paragraphs]
    print(f"INFO: docx has {len(paragraphs)} paragraphs, pptx has {num_slides} slides")

    # -----------------------------------------------------------------------
    # Component 1: All 14 'Slide N:' labels present in correct order (0.4 pts)
    # The docx must contain exactly 14 slide labels like 'Slide 1:', 'Slide 2:', ...
    # FAILS on initial_env (no docx exists — blocked by precondition gate above)
    # PASSES on golden_env (docx has all 14 labels in order)
    # -----------------------------------------------------------------------
    try:
        # Find all paragraphs matching 'Slide N:' pattern
        label_paras = []
        for para_text in paragraphs:
            m = re.match(r'^Slide\s+(\d+)\s*:', para_text.strip(), re.IGNORECASE)
            if m:
                label_paras.append(int(m.group(1)))

        expected_labels = list(range(1, num_slides + 1))

        if len(label_paras) == num_slides and label_paras == expected_labels:
            print(f"PASS: Component 1 — All {num_slides} 'Slide N:' labels present in correct order (0.4 pts)")
            total_score += 0.4
        elif len(label_paras) > 0:
            # Partial credit proportional to number of valid labels found
            fraction = len(set(label_paras) & set(expected_labels)) / num_slides
            partial = round(0.4 * fraction, 2)
            print(f"PARTIAL: Component 1 — Found {len(label_paras)}/{num_slides} slide labels (order/completeness issue): {partial} pts")
            if partial > 0:
                total_score += partial
        else:
            print(f"FAIL: Component 1 — No 'Slide N:' labels found in docx (0.0 pts)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # -----------------------------------------------------------------------
    # Component 2: Notes text content matches source pptx notes (0.4 pts)
    # Each slide's notes paragraph in docx should match corresponding pptx notes.
    # FAILS on initial_env (no docx — blocked by precondition gate)
    # PASSES on golden_env (all 14 slide notes match pptx source)
    # -----------------------------------------------------------------------
    try:
        # Build a map: slide_number -> docx paragraph text (paragraph after label)
        docx_notes_map = {}
        for i, para_text in enumerate(paragraphs):
            m = re.match(r'^Slide\s+(\d+)\s*:', para_text.strip(), re.IGNORECASE)
            if m:
                slide_num = int(m.group(1))
                # The notes paragraph should be immediately after the label
                if i + 1 < len(paragraphs):
                    docx_notes_map[slide_num] = paragraphs[i + 1]

        matched_slides = 0
        total_slides_checked = 0
        for slide_num, pptx_notes in slide_notes:
            if not pptx_notes:
                continue  # skip slides with no notes
            total_slides_checked += 1
            docx_text = docx_notes_map.get(slide_num, "")
            # Verify pptx notes content (first 80 chars) appears in docx paragraph
            pptx_notes_prefix = pptx_notes[:80].strip()
            if pptx_notes_prefix and pptx_notes_prefix in docx_text:
                matched_slides += 1
            else:
                print(f"  MISMATCH: Slide {slide_num} — "
                      f"expected '{pptx_notes_prefix[:50]}', got '{docx_text[:50]}'")

        if total_slides_checked > 0 and matched_slides == total_slides_checked:
            total_score += 0.4
            print(f"PASS: Component 2 — All {matched_slides}/{total_slides_checked} slides' notes content match pptx source (0.4 pts)")
        elif total_slides_checked > 0 and matched_slides > 0:
            comp2_score = round(0.4 * matched_slides / total_slides_checked, 2)
            total_score += comp2_score
            print(f"PARTIAL: Component 2 — {matched_slides}/{total_slides_checked} slides match pptx notes content ({comp2_score} pts)")
        else:
            print("FAIL: Component 2 — No slides with notes matched (0.0 pts)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # -----------------------------------------------------------------------
    # Component 3: Notes sections separated by blank paragraph lines (0.2 pts)
    # After each slide's notes, there should be a blank paragraph before the next label.
    # FAILS on initial_env (no docx — blocked by precondition gate)
    # PASSES on golden_env (all 13 inter-slide gaps have blank separator)
    # -----------------------------------------------------------------------
    try:
        separator_count = 0
        total_separators_expected = 0

        for i, para_text in enumerate(paragraphs):
            m = re.match(r'^Slide\s+(\d+)\s*:', para_text.strip(), re.IGNORECASE)
            if m:
                slide_num = int(m.group(1))
                # Only check inter-slide gaps (not after last slide)
                if slide_num < num_slides:
                    total_separators_expected += 1
                    # Para at i+1 is notes, Para at i+2 should be blank line
                    if i + 2 < len(paragraphs) and paragraphs[i + 2].strip() == "":
                        separator_count += 1
                    else:
                        next_text = paragraphs[i + 2] if i + 2 < len(paragraphs) else "(missing)"
                        print(f"  NO_BLANK: After Slide {slide_num} notes, expected blank, "
                              f"got: '{next_text[:50]}'")

        if total_separators_expected > 0 and separator_count == total_separators_expected:
            total_score += 0.2
            print(f"PASS: Component 3 — All {separator_count}/{total_separators_expected} blank line separators present (0.2 pts)")
        elif total_separators_expected > 0 and separator_count > 0:
            comp3_score = round(0.2 * separator_count / total_separators_expected, 2)
            total_score += comp3_score
            print(f"PARTIAL: Component 3 — {separator_count}/{total_separators_expected} blank separators present ({comp3_score} pts)")
        else:
            print("FAIL: Component 3 — No blank line separators found (0.0 pts)")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entrypoint: run verification
verify_task()
