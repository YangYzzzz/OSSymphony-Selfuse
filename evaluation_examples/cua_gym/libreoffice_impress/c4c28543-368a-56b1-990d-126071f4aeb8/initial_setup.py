"""
Initial Setup: Strategy presentation with presenter notes for notes export task
Task ID: osworld_multi_apps_impress_notes_export_009
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_impress_notes_export_009'
OUTPUT = f'{WORKDIR}/Strategy_2025.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


# Slide data: (title, subtitle_or_content, notes)
SLIDES_DATA = [
    {
        "layout": 0,  # Title Slide
        "title": "Acme Corporation",
        "subtitle": "Annual Strategy Review 2025",
        "notes": (
            "Welcome everyone to our annual strategy review. This year marks a pivotal moment for "
            "Acme Corporation as we navigate shifting market dynamics and accelerate our digital "
            "transformation agenda. We have a full agenda today, so let's get started."
        ),
    },
    {
        "layout": 1,  # Title + Content
        "title": "Agenda",
        "content": (
            "1. Executive Summary\n"
            "2. Market Analysis & Competitive Landscape\n"
            "3. 2024 Performance Review\n"
            "4. Strategic Priorities for 2025\n"
            "5. OKRs & Key Metrics\n"
            "6. Product Roadmap\n"
            "7. Financial Outlook\n"
            "8. Team & Talent Strategy\n"
            "9. Risk Management\n"
            "10. Q&A"
        ),
        "notes": (
            "Walk through the agenda briefly. Emphasize that we'll have dedicated Q&A at the end, "
            "but feel free to hold questions during each section as well. We're allocating about "
            "45 minutes for discussion after the main presentation. Each section head will present "
            "their own material."
        ),
    },
    {
        "layout": 1,
        "title": "Executive Summary",
        "content": (
            "• Revenue grew 18% YoY to $142M in 2024\n"
            "• Launched 3 major product lines ahead of schedule\n"
            "• Expanded into 4 new markets across APAC and EMEA\n"
            "• Customer NPS improved from 42 to 61\n"
            "• Headcount grew from 620 to 890 employees\n"
            "• 2025 target: $200M revenue, 35% gross margin"
        ),
        "notes": (
            "These headline numbers reflect a strong year. The 18% revenue growth is particularly "
            "notable given broader macroeconomic headwinds in H2 2024. Our NPS jump of 19 points "
            "is the result of the customer success initiative launched in Q1. The 2025 targets are "
            "ambitious but achievable based on our pipeline and market conditions. Highlight the "
            "headcount growth as evidence of our investment in execution capacity."
        ),
    },
    {
        "layout": 1,
        "title": "Market Analysis",
        "content": (
            "Total Addressable Market: $8.5B (growing 22% CAGR)\n\n"
            "Key Trends:\n"
            "• AI-driven workflow automation adoption accelerating\n"
            "• Enterprise SaaS consolidation continues\n"
            "• Shift toward outcome-based pricing models\n"
            "• Regulatory complexity increasing in EU and US\n\n"
            "Competitive Position: #3 in overall market share, #1 in SMB segment"
        ),
        "notes": (
            "The market analysis was conducted with input from Gartner and our own customer "
            "research team. The 22% CAGR is a revised upward estimate driven by faster-than-expected "
            "AI adoption. Our #3 position overall but #1 in SMB is a key strategic asset — we should "
            "not abandon our core SMB strength while pursuing enterprise growth. Be ready to discuss "
            "the competitive threats from both legacy vendors and newer entrants."
        ),
    },
    {
        "layout": 1,
        "title": "2024 Performance Review",
        "content": (
            "Q1: $31.2M revenue | 112% quota attainment\n"
            "Q2: $34.8M revenue | 118% quota attainment\n"
            "Q3: $36.5M revenue | 108% quota attainment\n"
            "Q4: $39.5M revenue | 122% quota attainment\n\n"
            "Churn Rate: 4.2% (down from 6.1% in 2023)\n"
            "Net Revenue Retention: 118%\n"
            "CAC Payback Period: 14 months"
        ),
        "notes": (
            "Q3 was our softest quarter due to the sales leadership transition and extended enterprise "
            "deal cycles. Recovery in Q4 was stronger than projected, driven by a focus on mid-market "
            "closures and a successful year-end promotion. Churn improvement from 6.1% to 4.2% is "
            "directly attributable to the onboarding revamp and customer health scoring model deployed "
            "in May 2024. NRR of 118% means our existing customers are expanding faster than they "
            "churn — this is a critical metric to sustain."
        ),
    },
    {
        "layout": 1,
        "title": "Strategic Priorities for 2025",
        "content": (
            "Priority 1: Accelerate Enterprise Market Penetration\n"
            "  — Dedicated enterprise sales team of 40 reps by Q2\n"
            "  — Launch enterprise-tier product by March 2025\n\n"
            "Priority 2: AI Product Differentiation\n"
            "  — Embed AI assistant across all product modules\n"
            "  — Release AI analytics dashboard in Q1\n\n"
            "Priority 3: International Expansion\n"
            "  — Open London and Singapore offices in H1\n"
            "  — Localize product for German and Japanese markets"
        ),
        "notes": (
            "These three priorities were selected after a rigorous strategic planning process "
            "involving all senior leadership and board input. Each priority has a dedicated owner: "
            "Elena Rodriguez for Enterprise, Priya Mehta for AI, and Daniel Kim for International. "
            "Stress that these are not just aspirations — each has defined quarterly milestones and "
            "budget allocations. Trade-offs were made: we are deprioritizing the consumer segment "
            "and the marketplace initiative to focus resources."
        ),
    },
    {
        "layout": 1,
        "title": "OKRs: Objective 1 — Enterprise Growth",
        "content": (
            "Objective: Become the preferred platform for mid-market and enterprise teams\n\n"
            "KR1: Close 80 new enterprise accounts (>$50K ACV) by Dec 2025\n"
            "KR2: Achieve $60M in new ARR from enterprise segment\n"
            "KR3: Reduce enterprise sales cycle from 90 to 65 days\n"
            "KR4: Reach 95% CSAT score among enterprise customers\n"
            "KR5: Launch 2 enterprise-specific integrations (SAP, Salesforce)"
        ),
        "notes": (
            "The enterprise objective is our highest-priority growth lever. The 80 new accounts "
            "target was set based on our current pipeline capacity and hiring plan. Each of the 40 "
            "enterprise reps will be expected to close 2 accounts per year on average. The SAP and "
            "Salesforce integrations are specifically requested by our top 3 enterprise prospects and "
            "will be a significant unlock. The 65-day sales cycle target requires improvements in "
            "our security review and procurement processes — Legal and IT have committed resources."
        ),
    },
    {
        "layout": 1,
        "title": "OKRs: Objective 2 — AI Leadership",
        "content": (
            "Objective: Deliver AI capabilities that measurably improve user productivity\n\n"
            "KR1: 60% of active users adopt AI features by Q3 2025\n"
            "KR2: AI features reduce average task completion time by 30%\n"
            "KR3: Launch 5 AI-powered automation workflows\n"
            "KR4: Achieve top-3 ranking in G2 AI features category\n"
            "KR5: File 3 patents related to our AI workflow engine"
        ),
        "notes": (
            "AI is both a product differentiator and a retention driver. The 60% adoption target "
            "for AI features is aggressive but achievable with the right in-app onboarding and "
            "education investment. The 30% task completion reduction is based on beta testing data "
            "from our early access program with 200 customers. The G2 ranking objective requires "
            "a coordinated marketing and review generation campaign starting in Q1. Our IP attorney "
            "has reviewed the 3 patent applications — they are on track for filing by Q2."
        ),
    },
    {
        "layout": 1,
        "title": "OKRs: Objective 3 — International Expansion",
        "content": (
            "Objective: Establish profitable operations in 4 new international markets\n\n"
            "KR1: Generate $15M ARR from international markets by Dec 2025\n"
            "KR2: Hire 45 employees across London and Singapore offices\n"
            "KR3: Achieve product localization for DE and JP by Q2\n"
            "KR4: Sign 3 channel partners in EMEA\n"
            "KR5: Obtain SOC 2 Type II certification by Q3"
        ),
        "notes": (
            "International expansion is our highest-risk strategic initiative and requires careful "
            "execution. The London office lease has been signed; Singapore is in final negotiations. "
            "We've hired Daniel Kim as VP International to lead this initiative — he brings 12 years "
            "of experience scaling SaaS companies in EMEA and APAC. The SOC 2 Type II certification "
            "is mandatory for enterprise deals in Europe and is already in progress with Deloitte. "
            "Channel partners in EMEA will be critical for market entry — we have 5 candidates in "
            "the pipeline currently."
        ),
    },
    {
        "layout": 1,
        "title": "Product Roadmap H1 2025",
        "content": (
            "January:\n"
            "  • AI Analytics Dashboard (GA release)\n"
            "  • Enterprise SSO & SCIM provisioning\n\n"
            "February:\n"
            "  • Advanced Reporting Suite v2\n"
            "  • Mobile app redesign (iOS & Android)\n\n"
            "March:\n"
            "  • Enterprise tier product launch\n"
            "  • Salesforce integration beta\n\n"
            "Q2 (April–June):\n"
            "  • SAP integration GA\n"
            "  • German & Japanese localization\n"
            "  • AI Workflow Automation v1"
        ),
        "notes": (
            "The H1 roadmap has been reviewed and approved by the Product Council. All items in "
            "January and February have dependencies already resolved. The Enterprise tier launch in "
            "March is the most critical — it's blocked on Legal approval of new contract templates "
            "and Infra completion of the multi-tenant isolation work. Priya and the engineering leads "
            "have a daily standup on these blockers. The mobile redesign in February incorporates "
            "extensive user research from 150 interviews conducted in Q4 2024."
        ),
    },
    {
        "layout": 1,
        "title": "Product Roadmap H2 2025",
        "content": (
            "Q3 (July–September):\n"
            "  • AI Workflow Automation v2 with custom rules\n"
            "  • Partner API marketplace launch\n"
            "  • Offline mode for mobile\n\n"
            "Q4 (October–December):\n"
            "  • Predictive analytics module\n"
            "  • Advanced permissions & governance\n"
            "  • Self-hosted / private cloud deployment option\n\n"
            "Strategic Bets:\n"
            "  • Voice interface for hands-free workflows\n"
            "  • Real-time collaboration v2"
        ),
        "notes": (
            "H2 priorities are subject to revision based on H1 execution and market feedback. "
            "The Partner API marketplace is a strategic initiative that could significantly expand "
            "our ecosystem but requires substantial platform work. The self-hosted deployment option "
            "is a direct response to 23 lost enterprise deals where data residency was the blocking "
            "factor. Voice interface and real-time collaboration v2 are on our strategic bets list — "
            "they won't be cut but timelines may flex. We'll do a mid-year roadmap review in June."
        ),
    },
    {
        "layout": 1,
        "title": "Financial Outlook 2025",
        "content": (
            "Revenue Target: $200M (+41% YoY)\n"
            "  • New ARR: $85M\n"
            "  • Expansion ARR: $28M\n"
            "  • Churn: -$13M\n\n"
            "Gross Margin: 35% (target 38% by Q4)\n"
            "Operating Expenses: $185M\n"
            "EBITDA Target: Breakeven by Q4 2025\n\n"
            "Key Investments:\n"
            "  • R&D: $62M (up from $44M)\n"
            "  • Sales & Marketing: $78M\n"
            "  • G&A: $22M\n"
            "  • International expansion: $23M"
        ),
        "notes": (
            "The financial model assumes 90% of new ARR coming from direct sales and 10% from "
            "channel partners. The gross margin improvement from current 32% to 38% by Q4 depends "
            "on infrastructure optimization work and the new cloud cost management program. Our CFO "
            "Rachel Torres has built three scenarios: base, upside (+15%), and downside (-20%). "
            "The board has approved the base case budget. We have 18 months of runway at current "
            "burn. Series D fundraising discussions are planned for Q3 if we hit our Q2 milestones."
        ),
    },
    {
        "layout": 1,
        "title": "Team & Talent Strategy",
        "content": (
            "Current headcount: 890 employees\n"
            "2025 hiring plan: +210 employees (total: 1,100)\n\n"
            "Key hires:\n"
            "  • Chief Revenue Officer (search in progress)\n"
            "  • VP Engineering (3 finalists)\n"
            "  • 40 Enterprise Sales Representatives\n"
            "  • 25 Engineers (AI/ML focus)\n"
            "  • 45 International hires\n\n"
            "Culture Initiatives:\n"
            "  • Manager effectiveness program (Q1)\n"
            "  • 15% compensation increase for IC engineers\n"
            "  • Flexible work policy update"
        ),
        "notes": (
            "Talent is our most important execution risk for 2025. The CRO search is our top priority "
            "— we have engaged Spencer Stuart and expect to close by end of January. The VP Engineering "
            "search has 3 strong finalists; offer is expected next week. The 15% compensation increase "
            "for IC engineers was approved by the board in response to above-market attrition in H2 2024. "
            "Our People Analytics team has identified manager effectiveness as the #1 driver of employee "
            "satisfaction — the manager program launching in Q1 is a direct response. We are also "
            "updating our flexible work policy to remain competitive in a tight labor market."
        ),
    },
    {
        "layout": 1,
        "title": "Risk Management & Mitigation",
        "content": (
            "Risk 1: Enterprise sales ramp slower than planned\n"
            "  Mitigation: Channel partner program accelerates pipeline\n\n"
            "Risk 2: AI feature adoption below target\n"
            "  Mitigation: In-app education, success team playbooks\n\n"
            "Risk 3: International execution complexity\n"
            "  Mitigation: Experienced VP International hire, phased approach\n\n"
            "Risk 4: Talent gaps in key roles\n"
            "  Mitigation: Aggressive compensation, contractor backup plan\n\n"
            "Risk 5: Macro deterioration reduces IT budgets\n"
            "  Mitigation: ROI-focused sales motion, flexible pricing"
        ),
        "notes": (
            "Risk management is a regular agenda item at the monthly leadership team meetings and "
            "quarterly board reviews. Each risk has a named owner and a defined escalation threshold. "
            "Risk 1 is our primary concern — if the enterprise ramp is 30% below plan by Q2, we will "
            "accelerate the channel partner program and consider an acquisition of a small enterprise "
            "sales team. Risk 5 (macro) is outside our control but our flexible pricing initiative "
            "and ROI calculator tool are designed to help customers justify spend in tighter budgets. "
            "We will review this risk register monthly and update the board on material changes."
        ),
    },
]


def create_initial():
    # Ensure Desktop directory exists
    os.makedirs(WORKDIR, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for slide_info in SLIDES_DATA:
        layout_idx = slide_info["layout"]
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = slide_info["title"]

        # Set content/subtitle
        if layout_idx == 0:
            # Title slide: placeholder[1] is subtitle
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_info.get("subtitle", "")
        else:
            # Title + Content: placeholder[1] is body content
            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.text = slide_info.get("content", "")

        # Add presenter notes
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = slide_info["notes"]

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open the presentation in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
