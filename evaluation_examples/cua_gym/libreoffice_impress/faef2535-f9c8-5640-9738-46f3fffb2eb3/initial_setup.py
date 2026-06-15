"""
Initial Setup: Create a 4-slide curriculum presentation with a flat bulleted list on slide 2.
Task ID: impstruct_035
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impstruct_035'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title Slide ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Data Analytics Curriculum 2025"
    slide1.placeholders[1].text = "Professional Development Program\nQuarter 3 - Advanced Track"

    # ── Slide 2: Course Outline (flat bulleted list — 9 items, all level 0) ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Course Outline"

    items = [
        "Introduction to Statistical Modeling",
        "Probability distributions and sampling methods",
        "Hypothesis testing with real-world datasets",
        "Machine Learning Fundamentals",
        "Supervised learning algorithms and evaluation metrics",
        "Feature engineering and cross-validation techniques",
        "Data Visualization Best Practices",
        "Dashboard design principles and color theory",
        "Interactive reporting with modern BI tools",
    ]

    body = slide2.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for idx, item in enumerate(items):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0  # all flat — no indentation
        p.font.size = Pt(18)

    # ── Slide 3: Schedule ──
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Program Schedule"

    schedule_items = [
        "Week 1-3: Statistical Foundations",
        "Week 4-6: Machine Learning Core",
        "Week 7-9: Visualization & Reporting",
        "Week 10: Capstone Project Presentations",
        "Total Duration: 10 weeks (40 hours)",
    ]

    body3 = slide3.placeholders[1]
    tf3 = body3.text_frame
    tf3.clear()

    for idx, item in enumerate(schedule_items):
        if idx == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)

    # ── Slide 4: Contact & Resources ──
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Contact & Resources"

    contact_items = [
        "Program Lead: Dr. Amara Osei (a.osei@analytics-academy.edu)",
        "Teaching Assistants: Wei Zhang, Priya Sharma",
        "Office Hours: Tuesdays & Thursdays 2:00 - 4:00 PM",
        "Online Portal: academy.analytics-hub.io/q3-advanced",
        "Required Textbook: Applied Predictive Analytics, 4th Edition",
    ]

    body4 = slide4.placeholders[1]
    tf4 = body4.text_frame
    tf4.clear()

    for idx, item in enumerate(contact_items):
        if idx == 0:
            p = tf4.paragraphs[0]
        else:
            p = tf4.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
