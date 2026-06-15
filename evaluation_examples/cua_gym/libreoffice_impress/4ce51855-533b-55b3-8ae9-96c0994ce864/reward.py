"""
Reward Script: Insert a new slide after slide 1 with title 'Course Outline' and 4 bulleted items.
Task ID: impress_teach_004
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.2): Slide count increased to 4
  - Component 2 (0.2): Slide 2 title is 'Course Outline'
  - Component 3 (0.4): Slide 2 content has 4 correct bullet items in order
  - Component 4 (0.2): Original slides preserved (shifted to positions 3 & 4)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_004'

EXPECTED_BULLETS = [
    'Week 1-3: Fundamentals',
    'Week 4-6: Advanced Concepts',
    'Week 7-9: Applications',
    'Week 10: Final Review',
]

# Original slide titles that should now be at positions 3 and 4
ORIGINAL_SHIFTED_TITLES = [
    'Why Study Chemistry?',
    'Key Topics This Semester',
]


def persist_app_state(domain):
    """Save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_slide_title(slide):
    """Extract title text from a slide."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name and 'Title' in shape.name:
            return shape.text_frame.text.strip()
    # Fallback: check if shape is a title placeholder
    if slide.shapes.title is not None:
        return slide.shapes.title.text.strip()
    return ""


def get_content_paragraphs(slide):
    """Extract non-empty paragraph texts from content placeholder."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name and 'Content' in shape.name:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return texts


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Slide count is 4 (0.2 points)
    # Initial has 3 slides; golden should have 4 (one new slide inserted)
    try:
        if num_slides == 4:
            print(f"PASS: Component 1 -- Slide count is 4 (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 1 -- Expected 4 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Slide 2 title is 'Course Outline' (0.2 points)
    # In initial, slide 2 has title 'Why Study Chemistry?' -- this check only passes
    # when the new slide has been inserted at position 2
    try:
        if num_slides >= 2:
            slide2_title = get_slide_title(prs.slides[1])
            if slide2_title == 'Course Outline':
                print(f"PASS: Component 2 -- Slide 2 title is 'Course Outline' (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 2 -- Expected slide 2 title 'Course Outline', found '{slide2_title}'")
        else:
            print(f"FAIL: Component 2 -- Not enough slides (need at least 2)")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Slide 2 has 4 correct bullet items in order (0.4 points)
    # Award 0.1 per correct bullet in correct position
    # Initial slide 2 has completely different content, so this scores 0.0 on initial
    try:
        if num_slides >= 2:
            content_texts = get_content_paragraphs(prs.slides[1])
            bullet_score = 0.0
            matched_count = 0
            for idx, expected in enumerate(EXPECTED_BULLETS):
                if idx < len(content_texts) and content_texts[idx] == expected:
                    matched_count += 1
                    bullet_score += 0.1
                    print(f"  PASS: Bullet {idx+1} matches: '{expected}'")
                elif idx < len(content_texts):
                    print(f"  FAIL: Bullet {idx+1} expected '{expected}', found '{content_texts[idx]}'")
                else:
                    print(f"  FAIL: Bullet {idx+1} missing (expected '{expected}')")

            if matched_count == 4:
                print(f"PASS: Component 3 -- All 4 bullets correct (0.4 pts)")
                total_score += 0.4
            elif matched_count > 0:
                print(f"PARTIAL: Component 3 -- {matched_count}/4 bullets correct ({bullet_score} pts)")
                total_score += bullet_score
            else:
                print(f"FAIL: Component 3 -- 0/4 bullets correct")
        else:
            print(f"FAIL: Component 3 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Original slides preserved at positions 3 & 4 (0.2 points)
    # In initial, these are at positions 2 & 3. In golden, they should be shifted to 3 & 4.
    # This check fails on initial because slide 3 (initial) has title 'Key Topics This Semester',
    # not 'Why Study Chemistry?'. And initial only has 3 slides total.
    try:
        if num_slides >= 4:
            slide3_title = get_slide_title(prs.slides[2])
            slide4_title = get_slide_title(prs.slides[3])
            if (slide3_title == ORIGINAL_SHIFTED_TITLES[0] and
                    slide4_title == ORIGINAL_SHIFTED_TITLES[1]):
                print(f"PASS: Component 4 -- Original slides preserved at positions 3 & 4 (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 4 -- Slide 3 title='{slide3_title}' (expected '{ORIGINAL_SHIFTED_TITLES[0]}'), "
                      f"Slide 4 title='{slide4_title}' (expected '{ORIGINAL_SHIFTED_TITLES[1]}')")
        else:
            print(f"FAIL: Component 4 -- Not enough slides (need 4, found {num_slides})")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
