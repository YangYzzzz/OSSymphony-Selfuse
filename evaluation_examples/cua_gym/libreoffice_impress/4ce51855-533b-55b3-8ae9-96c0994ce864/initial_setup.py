"""
Initial Setup: Chemistry Course presentation with 3 slides
Task ID: impress_teach_004
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_004'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Introduction to Chemistry"
    slide1.placeholders[1].text = "A Comprehensive Course for Undergraduate Students\nDr. Elena Rodriguez | Fall 2025"

    # --- Slide 2: Course Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Why Study Chemistry?"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Chemistry is the central science connecting physics, biology, and earth sciences"
    p2 = tf2.add_paragraph()
    p2.text = "Understanding molecular interactions helps explain everyday phenomena"
    p2.level = 0
    p3 = tf2.add_paragraph()
    p3.text = "Applications range from medicine to materials engineering"
    p3.level = 0
    p4 = tf2.add_paragraph()
    p4.text = "Critical thinking and problem-solving skills development"
    p4.level = 0

    # --- Slide 3: Key Topics ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Key Topics This Semester"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Atomic Structure and Periodic Trends"
    for topic in [
        "Chemical Bonding and Molecular Geometry",
        "Stoichiometry and Chemical Reactions",
        "Thermodynamics and Kinetics",
        "Acids, Bases, and Equilibrium",
        "Organic Chemistry Fundamentals",
    ]:
        p = tf3.add_paragraph()
        p.text = topic
        p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
