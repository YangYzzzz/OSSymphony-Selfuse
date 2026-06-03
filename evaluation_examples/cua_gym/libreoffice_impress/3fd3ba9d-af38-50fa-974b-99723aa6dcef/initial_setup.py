"""
Initial Setup: 10-slide product catalog with unstructured content on blank layout
Task ID: impress_gf2_045
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_045'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


# Product data for 10 slides - realistic catalog items
products = [
    {
        "name": "ErgoFlex Pro Standing Desk",
        "category": "Office Furniture",
        "price": "$749.99",
        "sku": "EFP-2025-001",
        "description": "Height-adjustable standing desk with memory presets and cable management.",
        "dimensions": "60\" x 30\" x 25-50\"",
        "weight": "85 lbs",
        "color": "Walnut / Matte Black",
        "rating": "4.8/5",
        "stock": "In Stock",
    },
    {
        "name": "AeroGlide Wireless Mouse",
        "category": "Computer Peripherals",
        "price": "$59.99",
        "sku": "AGM-2025-017",
        "description": "Ultra-lightweight wireless mouse with 20,000 DPI sensor and 80-hour battery.",
        "dimensions": "4.9\" x 2.5\" x 1.5\"",
        "weight": "63g",
        "color": "Pearl White",
        "rating": "4.6/5",
        "stock": "In Stock",
    },
    {
        "name": "SonicPure ANC Headphones",
        "category": "Audio Equipment",
        "price": "$329.00",
        "sku": "SPA-2025-042",
        "description": "Premium over-ear headphones with adaptive noise cancellation and spatial audio.",
        "dimensions": "7.5\" x 6.8\" x 3.4\"",
        "weight": "254g",
        "color": "Midnight Blue",
        "rating": "4.9/5",
        "stock": "Limited Stock",
    },
    {
        "name": "LumaView 4K Monitor",
        "category": "Displays",
        "price": "$899.00",
        "sku": "LV4-2025-008",
        "description": "32-inch 4K IPS display with 99% DCI-P3 color gamut and USB-C hub.",
        "dimensions": "28.2\" x 20.1\" x 9.3\"",
        "weight": "17.6 lbs",
        "color": "Space Gray",
        "rating": "4.7/5",
        "stock": "In Stock",
    },
    {
        "name": "TypeMaster Mechanical Keyboard",
        "category": "Computer Peripherals",
        "price": "$179.99",
        "sku": "TMK-2025-033",
        "description": "Hot-swappable mechanical keyboard with RGB per-key lighting and PBT keycaps.",
        "dimensions": "17.3\" x 5.1\" x 1.4\"",
        "weight": "2.1 lbs",
        "color": "Charcoal",
        "rating": "4.5/5",
        "stock": "In Stock",
    },
    {
        "name": "SwiftCharge 100W GaN Charger",
        "category": "Power & Charging",
        "price": "$69.99",
        "sku": "SC1-2025-051",
        "description": "Compact 100W GaN charger with 3 USB-C ports and 1 USB-A port.",
        "dimensions": "2.6\" x 2.6\" x 1.3\"",
        "weight": "7.2 oz",
        "color": "White",
        "rating": "4.8/5",
        "stock": "In Stock",
    },
    {
        "name": "CloudSync NAS Pro 4-Bay",
        "category": "Storage Solutions",
        "price": "$549.00",
        "sku": "CNP-2025-012",
        "description": "4-bay NAS enclosure with dual 2.5GbE, hardware RAID, and cloud backup.",
        "dimensions": "9.1\" x 6.3\" x 8.7\"",
        "weight": "5.9 lbs",
        "color": "Black",
        "rating": "4.4/5",
        "stock": "Pre-order",
    },
    {
        "name": "FlexMount Dual Monitor Arm",
        "category": "Office Accessories",
        "price": "$129.99",
        "sku": "FMA-2025-029",
        "description": "Gas spring dual monitor arm supporting up to 32\" displays, 360 rotation.",
        "dimensions": "Clamp: 0.4-3.5\" desk thickness",
        "weight": "11.2 lbs",
        "color": "Matte Silver",
        "rating": "4.7/5",
        "stock": "In Stock",
    },
    {
        "name": "PixelPen Pro Drawing Tablet",
        "category": "Creative Tools",
        "price": "$399.00",
        "sku": "PPT-2025-006",
        "description": "13.3\" pen display tablet with 8192 pressure levels and tilt recognition.",
        "dimensions": "15.2\" x 9.6\" x 0.5\"",
        "weight": "1.9 lbs",
        "color": "Graphite",
        "rating": "4.6/5",
        "stock": "In Stock",
    },
    {
        "name": "AirStream Desk Fan",
        "category": "Office Comfort",
        "price": "$44.99",
        "sku": "ADF-2025-078",
        "description": "Bladeless USB-C desk fan with 12-speed control and ambient noise under 25dB.",
        "dimensions": "6.3\" x 6.3\" x 11.4\"",
        "weight": "1.4 lbs",
        "color": "Arctic White",
        "rating": "4.3/5",
        "stock": "In Stock",
    },
]


def create_initial():
    prs = Presentation()
    # Default slide dimensions (widescreen 13.33" x 7.5")

    for i, prod in enumerate(products):
        # Use blank layout (index 6 = blank in default template)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add unstructured content - just text boxes scattered around
        # Product title - placed somewhat randomly
        title_box = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(20), Cm(2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = prod["name"]
        run = p.runs[0]
        run.font.size = Pt(28)
        run.font.bold = True

        # Category subtitle
        cat_box = slide.shapes.add_textbox(Cm(1), Cm(2.5), Cm(15), Cm(1.2))
        tf2 = cat_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"Category: {prod['category']}"
        run2 = p2.runs[0]
        run2.font.size = Pt(14)
        run2.font.italic = True
        run2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

        # Description - large text block in the middle
        desc_box = slide.shapes.add_textbox(Cm(2), Cm(4.5), Cm(22), Cm(3))
        tf3 = desc_box.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.text = prod["description"]
        run3 = p3.runs[0]
        run3.font.size = Pt(16)

        # Specs dumped as a single block - not in a table
        specs_box = slide.shapes.add_textbox(Cm(2), Cm(8), Cm(20), Cm(6))
        tf4 = specs_box.text_frame
        tf4.word_wrap = True
        specs_text = (
            f"SKU: {prod['sku']}  |  Price: {prod['price']}  |  "
            f"Dimensions: {prod['dimensions']}\n"
            f"Weight: {prod['weight']}  |  Color: {prod['color']}  |  "
            f"Rating: {prod['rating']}  |  Availability: {prod['stock']}"
        )
        p4 = tf4.paragraphs[0]
        p4.text = specs_text
        run4 = p4.runs[0]
        run4.font.size = Pt(12)

        # A simple colored rectangle as decoration (not the accent bar)
        rect = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Cm(0), Cm(18), Cm(33.87), Cm(0.3)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0x99, 0x99, 0x99)
        rect.line.fill.background()

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
