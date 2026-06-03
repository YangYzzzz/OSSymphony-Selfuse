"""
Reward Script: Product Catalog Presentation with Master Layout, Color Scheme, and Navigation
Task ID: impress_gf2_045
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): Master slide background is #FAFAFA
  Component 2 (0.15): Top accent bar on each slide (#2563EB, full width, ~0.5cm tall)
  Component 3 (0.15): Left product area placeholder on each slide (~40% width)
  Component 4 (0.20): Right-side table with product details on each slide
  Component 5 (0.30): Navigation arrows with correct hyperlinks on every slide
"""

import os
import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_045'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have exactly 10 slides
    num_slides = len(prs.slides)
    if num_slides != 10:
        print(f"PRECONDITION FAIL: Expected 10 slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    slide_width = prs.slide_width  # expected 9144000 EMU

    # Component 1: Master slide background is #FAFAFA (0.20 points)
    # This should FAIL on initial (no solid fill) and PASS on golden
    try:
        master = prs.slide_masters[0]
        master_fill = master.background.fill
        if master_fill.type is not None and master_fill.type == 1:  # SOLID fill
            color = str(master_fill.fore_color.rgb).upper()
            if color == "FAFAFA":
                print(f"PASS: Component 1 — Master background is #FAFAFA (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 1 — Master background color is #{color}, expected #FAFAFA")
        else:
            print(f"FAIL: Component 1 — Master background fill type is {master_fill.type}, expected SOLID (1)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Top accent bar on each slide — #2563EB, full width, approx 0.5cm (180000 EMU) tall (0.15 points)
    # Initial has a different bar at the bottom, not at top with accent color
    try:
        accent_bar_count = 0
        for si, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    # Check: positioned at top (top < 50000 EMU), full width (>= 80% slide width), short height
                    if shape.top < 50000 and shape.width >= slide_width * 0.8 and shape.height <= 400000:
                        try:
                            fill = shape.fill
                            if fill.type == 1:  # SOLID
                                color = str(fill.fore_color.rgb).upper()
                                if color == "2563EB":
                                    accent_bar_count += 1
                        except:
                            pass
        if accent_bar_count >= 10:
            print(f"PASS: Component 2 — All 10 slides have #2563EB accent bar at top (0.15 pts)")
            total_score += 0.15
        elif accent_bar_count >= 5:
            partial = 0.15 * (accent_bar_count / 10)
            print(f"PARTIAL: Component 2 — {accent_bar_count}/10 slides have accent bar ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — Only {accent_bar_count}/10 slides have #2563EB accent bar at top")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Left product area placeholder on each slide (0.15 points)
    # Golden has an auto shape at left side (~180000 left, ~540000 top, ~4500000 width, ~5400000 height)
    # with #E5E7EB border/fill color, representing ~40% width product image area.
    # Initial has NO such shape.
    try:
        left_area_count = 0
        for si, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    # Left-side large rectangle: left < 25% of slide, width roughly 35-50% of slide
                    shape_right = shape.left + shape.width
                    width_pct = shape.width / slide_width
                    if shape.left < slide_width * 0.15 and 0.30 <= width_pct <= 0.60 and shape.height > 3000000:
                        left_area_count += 1
                        break  # one per slide is enough
        if left_area_count >= 10:
            print(f"PASS: Component 3 — All 10 slides have left product area placeholder (0.15 pts)")
            total_score += 0.15
        elif left_area_count >= 5:
            partial = 0.15 * (left_area_count / 10)
            print(f"PARTIAL: Component 3 — {left_area_count}/10 slides have left area ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — Only {left_area_count}/10 slides have left product area")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Right-side table with product details on each slide (0.20 points)
    # Golden has TABLE shapes on the right side. Initial has NO tables.
    try:
        table_count = 0
        for si, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    # Table should be on right side (left > 40% of slide width)
                    if shape.left > slide_width * 0.40:
                        table = shape.table
                        # Must have at least 3 rows (some product details)
                        if len(table.rows) >= 3:
                            table_count += 1
                            break
        if table_count >= 10:
            print(f"PASS: Component 4 — All 10 slides have right-side product details table (0.20 pts)")
            total_score += 0.20
        elif table_count >= 5:
            partial = 0.20 * (table_count / 10)
            print(f"PARTIAL: Component 4 — {table_count}/10 slides have right-side table ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — Only {table_count}/10 slides have right-side product table")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Navigation arrows with correct hyperlinks (0.30 points)
    # Slide 1: only forward arrow (no back)
    # Slides 2-9: both forward and backward arrows
    # Slide 10: only backward arrow (no forward)
    # Arrows are AUTO_SHAPE with hyperlinks in the XML pointing to slideN.xml
    # Initial has NO hyperlinked arrow shapes.
    try:
        ns = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        }
        r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

        correct_nav = 0  # out of 10 slides

        with zipfile.ZipFile(file_path, 'r') as zf:
            for si in range(10):
                slide_num = si + 1
                slide_xml = f'ppt/slides/slide{slide_num}.xml'
                try:
                    with zf.open(slide_xml) as f:
                        root = ET.parse(f).getroot()

                    # Parse rels file for this slide to map rIds to targets
                    rels_path = f'ppt/slides/_rels/slide{slide_num}.xml.rels'
                    rid_map = {}
                    try:
                        with zf.open(rels_path) as rf:
                            rels_root = ET.parse(rf).getroot()
                            for rel in rels_root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                                rid = rel.get('Id')
                                target = rel.get('Target')
                                rid_map[rid] = target
                    except:
                        pass

                    # Find all shapes with hlinkClick
                    hyperlinked_targets = []
                    for sp in root.findall('.//p:sp', ns):
                        cNvPr = sp.find('.//p:nvSpPr/p:cNvPr', ns)
                        if cNvPr is not None:
                            hlk = cNvPr.find('a:hlinkClick', ns)
                            if hlk is not None:
                                rid = hlk.get(f'{{{r_ns}}}id', '')
                                target = rid_map.get(rid, '')
                                if target:
                                    hyperlinked_targets.append(target)

                    # Determine expected links
                    expected_forward = f'slide{slide_num + 1}.xml' if slide_num < 10 else None
                    expected_backward = f'slide{slide_num - 1}.xml' if slide_num > 1 else None

                    has_forward = any(expected_forward in t for t in hyperlinked_targets) if expected_forward else True
                    has_backward = any(expected_backward in t for t in hyperlinked_targets) if expected_backward else True

                    # Slide 1: no backward needed, just forward
                    # Slide 10: no forward needed, just backward
                    if has_forward and has_backward:
                        correct_nav += 1
                    else:
                        missing = []
                        if not has_forward and expected_forward:
                            missing.append(f"forward to slide{slide_num+1}")
                        if not has_backward and expected_backward:
                            missing.append(f"backward to slide{slide_num-1}")
                        print(f"  Slide {slide_num}: missing {', '.join(missing)}")
                except Exception as e:
                    print(f"  Slide {slide_num}: XML parse error: {e}")

        if correct_nav >= 10:
            print(f"PASS: Component 5 — All 10 slides have correct navigation arrows (0.30 pts)")
            total_score += 0.30
        elif correct_nav >= 5:
            partial = 0.30 * (correct_nav / 10)
            print(f"PARTIAL: Component 5 — {correct_nav}/10 slides have correct navigation ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 — Only {correct_nav}/10 slides have correct navigation arrows")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice
def persist_app_state(domain):
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
