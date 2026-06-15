"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 131 is supposed to kick off the appendix section, so I’d like a text box that literally just says "Appendix" sitting dead-center along the bottom edge of that slide. How do I add it in LibreOffice Impress?
Generated: 2025-09-10 15:30:49
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation


def verify_appendix_slide(file_path: str) -> float:
    """Verify that slide 131 contains a centred bottom-edge text box reading 'Appendix'.

    Scoring (progressive):
        0.0 – file missing / unreadable / slide <131
        +0.6 – the word 'Appendix' (case-insensitive) appears on slide 131
        +0.4 – that shape is horizontally centred (±5 % width) AND sits near the bottom edge (≥80 % slide height)
        Returns a float in [0,1].
    """

    max_score = 1.0
    score = 0.0

    # ---------- Load presentation ----------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0  # nothing else to check

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Could not open presentation: {e}")
        return 0.0

    # ---------- Ensure slide 131 exists (index 130) ----------
    if len(prs.slides) < 131:
        print(f"✗ Presentation has only {len(prs.slides)} slides – need at least 131")
        return 0.0

    slide = prs.slides[130]  # zero-based index

    # Slide dimensions (EMU)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    print(f"Slide dimensions (EMU) – width: {slide_width}, height: {slide_height}")

    # ---------- Locate shapes containing the word 'Appendix' ----------
    appendix_shapes = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            if "appendix" in shape.text.strip().lower():
                appendix_shapes.append(shape)
                print(
                    f"✓ Found shape with text '{shape.text.strip()}' at left={shape.left}, top={shape.top}, width={shape.width}, height={shape.height}"
                )

    if appendix_shapes:
        score += 0.6  # Text requirement satisfied
        print("✓ 'Appendix' text present on slide 131 (0.6 points)")
    else:
        print("✗ No shape containing the word 'Appendix' found on slide 131")
        print(f"Total score: {score}/{max_score}")
        print(f"REWARD: {score}")
        return score  # cannot continue without the text

    # ---------- Verify position (centre-bottom) ----------
    horizontally_centred = False
    bottom_aligned = False

    tol_x = slide_width * 0.05  # 5 % horizontal tolerance
    bottom_threshold = slide_height * 0.80  # bottom 20 % of slide

    for shp in appendix_shapes:
        centre_x = shp.left + shp.width / 2
        delta_x = abs(centre_x - slide_width / 2)
        bottom_edge = shp.top + shp.height

        horiz_ok = delta_x <= tol_x
        bottom_ok = bottom_edge >= bottom_threshold

        print(
            f"  Shape centre-x delta: {delta_x}, bottom edge: {bottom_edge} (horiz_ok={horiz_ok}, bottom_ok={bottom_ok})"
        )

        if horiz_ok and bottom_ok:
            horizontally_centred = True
            bottom_aligned = True
            break  # one valid shape is enough

    if horizontally_centred and bottom_aligned:
        score += 0.4
        print("✓ Shape is centred horizontally and aligned near bottom (0.4 points)")
    else:
        print("✗ 'Appendix' shape not correctly centred/bottom-aligned")

    final_score = min(score, max_score)
    print(f"Total score: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score


# ----------------- SCRIPT ENTRY POINT -----------------
if __name__ == "__main__":
    TEST_FILE = "/home/user/slide_131_is_supposed_to_kick_off_the_appendix_section_so_id_like_a_text_box_that_literally_just_say_golden.pptx"
    verify_appendix_slide(TEST_FILE)

