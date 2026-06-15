"""
Initial Setup: Create an 8-slide presentation with various transitions.
Slide 2 has Fade at 1.5s. Slides 4 and 6 have no transition.
Task ID: impress_tm_040
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import zipfile
import shutil
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_040'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

# Namespaces for transition XML
NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main',
}


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_transition_to_slide_xml(zip_path, slide_num, transition_type, duration_ms):
    """
    Add a transition element to a slide XML inside the pptx zip.
    slide_num is 1-based. transition_type: 'fade', 'push', 'wipe', 'dissolve', etc.
    duration_ms: transition duration in milliseconds (e.g., 1500 for 1.5s).
    """
    tmp_path = zip_path + '.tmp'
    slide_xml_name = f'ppt/slides/slide{slide_num}.xml'

    with zipfile.ZipFile(zip_path, 'r') as zin, zipfile.ZipFile(tmp_path, 'w') as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == slide_xml_name:
                # Parse and modify
                root = ET.fromstring(data)
                ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'

                # Remove existing transition if any
                for existing in root.findall(f'{{{ns_p}}}transition'):
                    root.remove(existing)

                # Create transition element
                # <p:transition spd="slow" advTm="0">
                #   <p:fade />
                # </p:transition>
                # Duration is set via the 'dur' attribute (in ms) on the transition child
                # or via spd attribute. We'll use advClick="1" and set dur on the child.
                trans_elem = ET.SubElement(root, f'{{{ns_p}}}transition')
                trans_elem.set('spd', 'med')
                # advClick: advance on mouse click
                trans_elem.set('advClick', '1')

                # Add the specific transition type child
                child = ET.SubElement(trans_elem, f'{{{ns_p}}}{transition_type}')

                # Set duration via p14:dur attribute on the transition element
                # The standard way to set duration in ms is the 'dur' attribute
                # on the p:transition element itself (in ms string)
                trans_elem.set('dur', str(duration_ms))

                # Move transition before p:timing if it exists (schema order matters)
                # p:transition should come after p:cSld and before p:timing
                timing = root.find(f'{{{ns_p}}}timing')
                if timing is not None:
                    root.remove(trans_elem)
                    idx = list(root).index(timing)
                    root.insert(idx, trans_elem)

                data = ET.tostring(root, xml_declaration=True, encoding='UTF-8')
            zout.writestr(item, data)

    shutil.move(tmp_path, zip_path)


def create_initial():
    prs = Presentation()

    # Slide content data
    slide_contents = [
        {
            'title': 'Annual Strategy Review',
            'body': 'Fiscal Year 2025-2026\nPrepared by the Executive Planning Team\nConfidential',
            'layout': 0,
        },
        {
            'title': 'Market Overview',
            'body': '• Global market grew 12.3% YoY\n• APAC region leads with 18.7% growth\n• North America stable at 8.2%\n• European markets recovering at 6.5%',
            'layout': 1,
        },
        {
            'title': 'Revenue Breakdown',
            'body': '• Product Sales: $45.2M (52%)\n• Services: $28.7M (33%)\n• Licensing: $13.1M (15%)\n• Total Revenue: $87.0M',
            'layout': 1,
        },
        {
            'title': 'Growth Initiatives',
            'body': '• Expand into Southeast Asian markets\n• Launch premium tier by Q3 2026\n• Strategic partnership with TechCorp\n• Increase R&D budget by 20%',
            'layout': 1,
        },
        {
            'title': 'Team Performance',
            'body': '• Engineering: 94% sprint completion rate\n• Sales: 108% quota achievement\n• Customer Success: NPS score 72\n• Marketing: 2.3M qualified leads generated',
            'layout': 1,
        },
        {
            'title': 'Risk Assessment',
            'body': '• Supply chain disruption (Medium)\n• Regulatory changes in EU (High)\n• Currency fluctuation exposure (Low)\n• Talent retention challenges (Medium)',
            'layout': 1,
        },
        {
            'title': 'Financial Projections',
            'body': '• Q1 2026: $24.5M projected\n• Q2 2026: $26.8M projected\n• Q3 2026: $29.1M projected\n• Q4 2026: $31.6M projected',
            'layout': 1,
        },
        {
            'title': 'Next Steps & Action Items',
            'body': '• Board presentation: April 15, 2026\n• Budget finalization: April 30, 2026\n• Quarterly review: June 15, 2026\n• Strategy offsite: August 2026',
            'layout': 1,
        },
    ]

    for i, content in enumerate(slide_contents):
        layout_idx = content['layout']
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        slide.shapes.title.text = content['title']
        if layout_idx == 0:
            slide.placeholders[1].text = content['body']
        elif layout_idx == 1:
            slide.placeholders[1].text = content['body']

    prs.save(OUTPUT)
    print(f'Base presentation saved: {OUTPUT}')

    # Now add transitions via XML manipulation
    # Transition config: (slide_num_1based, type, duration_ms)
    # Slide 1: push at 2.0s
    # Slide 2: fade at 1.5s (the source)
    # Slide 3: wipe at 1.0s
    # Slide 4: NO transition
    # Slide 5: dissolve at 2.5s
    # Slide 6: NO transition
    # Slide 7: push at 1.0s
    # Slide 8: wipe at 0.5s
    transitions = [
        (1, 'push', 2000),
        (2, 'fade', 1500),
        (3, 'wipe', 1000),
        # slide 4: no transition
        (5, 'dissolve', 2500),
        # slide 6: no transition
        (7, 'push', 1000),
        (8, 'wipe', 500),
    ]

    for slide_num, trans_type, dur_ms in transitions:
        add_transition_to_slide_xml(OUTPUT, slide_num, trans_type, dur_ms)
        print(f'  Slide {slide_num}: {trans_type} at {dur_ms}ms')

    print(f'Initial file created with transitions: {OUTPUT}')

    # Launch LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
