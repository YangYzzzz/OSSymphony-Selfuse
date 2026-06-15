"""
Reward Script: Set background color of all slides to light cream (#FFF8E7)
Task ID: impress_teach_005
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): At least half of slides have #FFF8E7 background (partial credit)
  Component 2 (0.6): All 10 slides have #FFF8E7 background (full completion)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_005'
TARGET_COLOR = 'FFF8E7'
EXPECTED_SLIDE_COUNT = 10


def get_slide_bg_color(slide):
    """Get the background fill color of a slide as a hex string, or None."""
    try:
        fill = slide.background.fill
        if fill.type is not None and fill.type == 1:  # SOLID fill
            return str(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file must have exactly 10 slides
    num_slides = len(prs.slides)
    if num_slides != EXPECTED_SLIDE_COUNT:
        print(f"PRECONDITION FAIL: Expected {EXPECTED_SLIDE_COUNT} slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    # Count how many slides have the target background color
    matching_slides = 0
    for i, slide in enumerate(prs.slides):
        bg_color = get_slide_bg_color(slide)
        if bg_color == TARGET_COLOR:
            matching_slides += 1
            print(f"  Slide {i+1}: background={bg_color} -- MATCHES target")
        else:
            print(f"  Slide {i+1}: background={bg_color} -- does NOT match target {TARGET_COLOR}")

    print(f"\nMatching slides: {matching_slides}/{EXPECTED_SLIDE_COUNT}")

    # Component 1: At least half of slides have correct background (0.4 points)
    # This checks partial completion -- the agent changed at least some slides
    try:
        if matching_slides >= EXPECTED_SLIDE_COUNT // 2:
            print(f"PASS: Component 1 -- {matching_slides} slides have target bg (>= {EXPECTED_SLIDE_COUNT // 2}) (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 -- only {matching_slides} slides have target bg (need >= {EXPECTED_SLIDE_COUNT // 2})")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: ALL slides have the correct background (0.6 points)
    # This checks full completion -- every single slide was changed
    try:
        if matching_slides == EXPECTED_SLIDE_COUNT:
            print(f"PASS: Component 2 -- all {EXPECTED_SLIDE_COUNT} slides have target bg (0.6 pts)")
            total_score += 0.6
        else:
            print(f"FAIL: Component 2 -- {matching_slides}/{EXPECTED_SLIDE_COUNT} slides have target bg (need all {EXPECTED_SLIDE_COUNT})")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
