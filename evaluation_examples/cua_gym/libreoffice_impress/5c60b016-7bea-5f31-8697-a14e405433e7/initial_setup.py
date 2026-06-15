"""
Initial Setup: English Literature lecture presentation with white backgrounds
Task ID: impress_teach_005
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_005'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    # Set white background explicitly
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    # Set white background explicitly
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def create_initial():
    prs = Presentation()

    # Slide 1: Title Slide
    add_title_slide(
        prs,
        "English Literature: Romantic Era to Modernism",
        "Professor Elena Vasquez — Spring 2026 — Lecture 7"
    )

    # Slide 2: Course Overview
    add_content_slide(prs, "Today's Agenda", [
        "Review of the Romantic poets and their legacy",
        "Transition from Victorian realism to early Modernism",
        "Close reading: Keats vs. Eliot imagery techniques",
        "Group discussion: Nature as symbol across periods",
        "Assignment briefing: Comparative essay due March 28"
    ])

    # Slide 3: Romantic Poetry Recap
    add_content_slide(prs, "The Romantic Poets — Key Themes", [
        "Emotion and imagination over reason and order",
        "Reverence for nature as spiritual force",
        "The individual as hero: Byron, Shelley, Keats",
        "Rejection of industrial society",
        "Lyrical Ballads (1798) as the movement's manifesto"
    ])

    # Slide 4: Keats Focus
    add_content_slide(prs, "John Keats: Ode to a Nightingale (1819)", [
        "Explores the tension between mortality and art",
        "Sensory imagery: 'sunburnt mirth' and 'purple-stained mouth'",
        "The nightingale as symbol of eternal beauty",
        "Negative capability: embracing uncertainty",
        "Compare with Shelley's 'To a Skylark' — similar bird motif"
    ])

    # Slide 5: Victorian Transition
    add_content_slide(prs, "From Romanticism to Victorian Realism", [
        "Industrial Revolution reshapes literary concerns",
        "Charles Dickens: social criticism through narrative",
        "George Eliot: psychological depth in Middlemarch (1871)",
        "The dramatic monologue: Browning's 'My Last Duchess'",
        "Tennyson's In Memoriam: grief meets scientific doubt"
    ])

    # Slide 6: Early Modernism
    add_content_slide(prs, "The Modernist Break (1900–1930)", [
        "World War I shatters Enlightenment optimism",
        "Stream of consciousness: Woolf, Joyce, Faulkner",
        "Fragmentation of form mirrors fractured experience",
        "Ezra Pound: 'Make it new' — the Imagist manifesto",
        "T.S. Eliot's The Waste Land (1922) as cultural diagnosis"
    ])

    # Slide 7: Close Reading Exercise
    add_content_slide(prs, "Close Reading: Imagery Comparison", [
        "Keats: 'Season of mists and mellow fruitfulness'",
        "Eliot: 'April is the cruellest month, breeding / Lilacs'",
        "Both use seasonal imagery — but to opposite ends",
        "Keats celebrates natural abundance and ripeness",
        "Eliot subverts spring as painful, ironic renewal"
    ])

    # Slide 8: Discussion Questions
    add_content_slide(prs, "Group Discussion Prompts", [
        "How does the concept of 'nature' shift from Keats to Eliot?",
        "Is Modernism a rejection or an evolution of Romanticism?",
        "What role does war play in literary transformation?",
        "Can we find Romantic echoes in 21st-century poetry?",
        "How do these poets define the purpose of art?"
    ])

    # Slide 9: Recommended Readings
    add_content_slide(prs, "Further Reading for Next Week", [
        "Virginia Woolf — Mrs Dalloway (1925), chapters 1–4",
        "James Joyce — Dubliners, 'The Dead' (1914)",
        "W.B. Yeats — 'The Second Coming' (1920)",
        "Harold Bloom — The Western Canon, ch. 12",
        "Course reader: 'Modernism and Its Discontents' (pp. 87–112)"
    ])

    # Slide 10: Assignment Info
    add_content_slide(prs, "Upcoming Assignment", [
        "Comparative Essay: Romantic vs. Modernist Aesthetics",
        "Choose one Romantic and one Modernist poem",
        "Analyze imagery, tone, and thematic concerns",
        "1500–2000 words, MLA format, due March 28",
        "Office hours: Tuesdays 2–4 PM, Humanities Building 305"
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
