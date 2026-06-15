"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 293 of my LibreOffice Impress deck, I need to insert a 3×4 table and set every gridline to precisely 0.5 pt in #808080 gray. How do I do that?
Generated: 2025-09-10 19:29:54
Status: success
Model: azure-o3
Total Steps: 2
"""

from pptx import Presentation
import os
from pptx.oxml.ns import qn


def verify_table_task(
    file_path: str,
    slide_number: int = 293,
    expected_rows: int = 3,
    expected_cols: int = 4,
    expected_border_width_pt: float = 0.5,
    expected_color_hex: str = "808080",
):
    """Verify that slide `slide_number` contains a table of
    `expected_rows` × `expected_cols` whose every grid-line is exactly
    `expected_border_width_pt` points wide and coloured `expected_color_hex`.

    Returns a progressive score between 0.0 and 1.0.
    """
    print(f"Verifying presentation: {file_path}")

    total_score = 0.0  # progressive score
    max_score = 1.0

    # ------------------------------------------------------------------
    # 1 ) Load the presentation (no points – prerequisite only)
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File does not exist")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded presentation with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # ------------------------------------------------------------------
    # 2 ) Slide existence (0.2 points)
    # ------------------------------------------------------------------
    if slide_number - 1 >= len(prs.slides):
        print(f"✗ Slide {slide_number} not found")
        print(f"REWARD: {total_score}")
        return total_score

    slide = prs.slides[slide_number - 1]
    print(f"✓ Slide {slide_number} exists")
    total_score += 0.2

    # ------------------------------------------------------------------
    # 3 ) Locate a 3 × 4 table on the slide (0.3 points)
    # ------------------------------------------------------------------
    table_shape = None
    for shape in slide.shapes:
        if shape.has_table:
            tbl = shape.table
            if len(tbl.rows) == expected_rows and len(tbl.columns) == expected_cols:
                table_shape = shape
                break

    if table_shape is None:
        print("✗ Table with expected dimensions not found on the slide")
        print(f"REWARD: {total_score}")
        return total_score

    print("✓ Found table with expected dimensions (3×4)")
    total_score += 0.3

    # ------------------------------------------------------------------
    # 4 ) Verify every border – width & colour (0.25 each)
    # ------------------------------------------------------------------
    tbl = table_shape.table

    expected_width_emu = int(round(expected_border_width_pt * 12700))  # 1 pt = 12700 EMU
    width_tolerance = 1270  # ±0.1 pt tolerance in EMU

    width_ok = True
    color_ok = True

    borders_checked = 0
    width_failures = 0
    color_failures = 0

    for r in range(len(tbl.rows)):
        for c in range(len(tbl.columns)):
            cell = tbl.cell(r, c)
            tc = cell._tc
            tcPr = tc.find(qn("a:tcPr"))
            if tcPr is None:
                width_ok = False
                color_ok = False
                continue

            # Four sides: left, right, top, bottom
            for side in ("lnL", "lnR", "lnT", "lnB"):
                ln = tcPr.find(qn(f"a:{side}"))
                if ln is None:
                    # Missing line definition counts as failure
                    width_ok = False
                    color_ok = False
                    borders_checked += 1
                    width_failures += 1
                    color_failures += 1
                    continue

                borders_checked += 1

                # ----- Width check -----
                w_attr = ln.get("w")
                if w_attr is None:
                    width_ok = False
                    width_failures += 1
                else:
                    try:
                        w_val = int(w_attr)
                        if abs(w_val - expected_width_emu) > width_tolerance:
                            width_ok = False
                            width_failures += 1
                    except ValueError:
                        width_ok = False
                        width_failures += 1

                # ----- Colour check -----
                srgb = ln.find(".//" + qn("a:srgbClr"))
                if (
                    srgb is None
                    or srgb.get("val") is None
                    or srgb.get("val").lower() != expected_color_hex.lower()
                ):
                    color_ok = False
                    color_failures += 1

    if borders_checked == 0:
        print("✗ No borders were found to check – unexpected structure")
        print(f"REWARD: {total_score}")
        return total_score

    # Width scoring
    if width_ok:
        print("✓ All borders have the correct width (0.5 pt)")
        total_score += 0.25
    else:
        print(
            f"✗ Border width incorrect in {width_failures}/{borders_checked} checked borders"
        )

    # Colour scoring
    if color_ok:
        print("✓ All borders are coloured #808080")
        total_score += 0.25
    else:
        print(
            f"✗ Border colour incorrect in {color_failures}/{borders_checked} checked borders"
        )

    # ------------------------------------------------------------------
    # 5 ) Final score
    # ------------------------------------------------------------------
    final_score = min(total_score, max_score)
    print(f"Total score: {final_score} / {max_score}")
    print(f"REWARD: {final_score}")
    return final_score


# ----------------------------------------------------------------------
# Execute verification when the script is run directly
# ----------------------------------------------------------------------
if __name__ == "__main__":
    FILE_PATH = (
        "/home/user/"
        "on_slide_293_of_my_libreoffice_impress_deck_i_need_to_insert_a_34_table_and_set_every_gridline_to_pr_golden.pptx"
    )
    verify_table_task(FILE_PATH)
