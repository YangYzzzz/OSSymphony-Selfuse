"""
Initial Setup: Create a 14-slide presentation with no watermark.
Task ID: impress_fix_071
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_071'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a text box with a single styled paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    # Standard widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Use blank layout; fall back if index doesn't exist
    num_layouts = len(prs.slide_layouts)
    blank_layout = prs.slide_layouts[6] if num_layouts > 6 else prs.slide_layouts[num_layouts - 1]

    def add_title(slide, text):
        """Add title to a slide, using placeholder if available, else textbox."""
        if slide.shapes.title is not None:
            slide.shapes.title.text = text
        else:
            add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                        text, font_size=28, bold=True, alignment=PP_ALIGN.LEFT)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(slide1, "Meridian Technologies - Q4 Strategy Review")
    if 1 in slide1.placeholders:
        slide1.placeholders[1].text = "Prepared by Strategic Planning Division\nNovember 2025"
    else:
        add_textbox(slide1, Inches(2), Inches(4), Inches(9), Inches(1.5),
                    "Prepared by Strategic Planning Division\nNovember 2025",
                    font_size=18, alignment=PP_ALIGN.CENTER)

    # ---- Slide 2: Executive Summary ----
    slide2 = prs.slides.add_slide(blank_layout)
    add_title(slide2, "Executive Summary")
    add_textbox(slide2, Inches(0.8), Inches(1.8), Inches(11), Inches(5),
                "Meridian Technologies achieved 23% year-over-year revenue growth in Q3, "
                "driven by strong performance in our cloud infrastructure and AI services "
                "divisions. However, enterprise hardware sales declined 8% as customers "
                "shifted toward managed service models. Our strategic pivot toward "
                "subscription-based offerings is progressing ahead of schedule, with ARR "
                "now representing 62% of total revenue versus the 55% target.",
                font_size=16)

    # ---- Slide 3: Revenue Breakdown ----
    slide3 = prs.slides.add_slide(blank_layout)
    add_title(slide3, "Revenue Breakdown by Division")
    rows, cols = 6, 4
    tbl = slide3.shapes.add_table(rows, cols, Inches(1.5), Inches(2), Inches(10), Inches(3.5)).table
    headers = ["Division", "Q3 Revenue ($M)", "Q2 Revenue ($M)", "Growth (%)"]
    for c, h in enumerate(headers):
        tbl.cell(0, c).text = h
    data = [
        ["Cloud Infrastructure", "$187.4", "$162.1", "+15.6%"],
        ["AI & ML Services", "$134.8", "$98.3", "+37.1%"],
        ["Enterprise Hardware", "$89.2", "$97.0", "-8.0%"],
        ["Managed Services", "$112.6", "$95.7", "+17.7%"],
        ["Professional Consulting", "$56.3", "$51.2", "+10.0%"],
    ]
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            tbl.cell(r, c).text = val

    # ---- Slide 4: Market Position ----
    slide4 = prs.slides.add_slide(blank_layout)
    add_title(slide4, "Competitive Market Position")
    add_textbox(slide4, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5),
                "Key Competitive Advantages:\n\n"
                "1. Proprietary NexusAI platform with 99.97% uptime\n"
                "2. Integrated hybrid cloud solution spanning 14 regions\n"
                "3. Industry-leading customer retention rate of 94.2%\n"
                "4. Patent portfolio of 340+ active patents\n"
                "5. Strategic partnerships with 8 Fortune 100 companies",
                font_size=14)
    add_textbox(slide4, Inches(7), Inches(1.8), Inches(5.5), Inches(5),
                "Market Share by Segment:\n\n"
                "Cloud Infrastructure: 12.3% (up from 10.8%)\n"
                "AI Services: 8.7% (up from 5.2%)\n"
                "Managed Services: 6.1% (stable)\n"
                "Enterprise Hardware: 4.8% (down from 5.5%)",
                font_size=14)

    # ---- Slide 5: Customer Acquisition ----
    slide5 = prs.slides.add_slide(blank_layout)
    add_title(slide5, "Customer Acquisition & Retention")
    rows2, cols2 = 5, 3
    tbl2 = slide5.shapes.add_table(rows2, cols2, Inches(2), Inches(2.2), Inches(9), Inches(3)).table
    for c, h in enumerate(["Metric", "Q3 2025", "Target"]):
        tbl2.cell(0, c).text = h
    metrics = [
        ["New Enterprise Clients", "47", "40"],
        ["Net Revenue Retention", "118%", "115%"],
        ["Customer Satisfaction (NPS)", "72", "70"],
        ["Average Contract Value", "$2.4M", "$2.1M"],
    ]
    for r, row_data in enumerate(metrics, 1):
        for c, val in enumerate(row_data):
            tbl2.cell(r, c).text = val

    # ---- Slide 6: Product Roadmap ----
    slide6 = prs.slides.add_slide(blank_layout)
    add_title(slide6, "Product Roadmap - Next 12 Months")
    add_textbox(slide6, Inches(0.8), Inches(1.8), Inches(11), Inches(5),
                "Q4 2025: Launch NexusAI 3.0 with multimodal capabilities\n"
                "Q1 2026: Release CloudBridge hybrid integration suite\n"
                "Q2 2026: Deploy edge computing nodes in 6 new regions\n"
                "Q3 2026: GA of SecureVault zero-trust security platform\n\n"
                "Investment Allocation:\n"
                "  R&D: $78M (34% of budget)\n"
                "  Infrastructure: $52M (23%)\n"
                "  Sales & Marketing: $45M (20%)\n"
                "  Operations: $38M (17%)\n"
                "  Administrative: $14M (6%)",
                font_size=14)

    # ---- Slide 7: Financial Highlights ----
    slide7 = prs.slides.add_slide(blank_layout)
    add_title(slide7, "Financial Highlights")
    add_textbox(slide7, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5),
                "Total Revenue: $580.3M\nGross Margin: 67.4%\nOperating Income: $98.7M\n"
                "Free Cash Flow: $112.3M",
                font_size=16, bold=True)
    add_textbox(slide7, Inches(7), Inches(1.8), Inches(5.5), Inches(2.5),
                "EPS (diluted): $1.87\nCash & Equivalents: $2.1B\n"
                "Debt-to-Equity: 0.34\nR&D Spend: $78M",
                font_size=16, bold=True)

    # ---- Slide 8: Regional Performance ----
    slide8 = prs.slides.add_slide(blank_layout)
    add_title(slide8, "Regional Performance Overview")
    rows3, cols3 = 5, 4
    tbl3 = slide8.shapes.add_table(rows3, cols3, Inches(1.5), Inches(2.2), Inches(10), Inches(3)).table
    for c, h in enumerate(["Region", "Revenue ($M)", "Growth (%)", "Key Driver"]):
        tbl3.cell(0, c).text = h
    regions = [
        ["North America", "$312.8", "+18.2%", "Cloud adoption surge"],
        ["Europe", "$148.5", "+22.7%", "GDPR compliance demand"],
        ["Asia-Pacific", "$87.3", "+41.3%", "AI services expansion"],
        ["Rest of World", "$31.7", "+12.1%", "Managed services"],
    ]
    for r, row_data in enumerate(regions, 1):
        for c, val in enumerate(row_data):
            tbl3.cell(r, c).text = val

    # ---- Slide 9: Team & Hiring ----
    slide9 = prs.slides.add_slide(blank_layout)
    add_title(slide9, "Team Growth & Talent Strategy")
    add_textbox(slide9, Inches(0.8), Inches(1.8), Inches(11), Inches(5),
                "Current Headcount: 4,827 (up 14% YoY)\n\n"
                "Key Hiring Areas:\n"
                "  - AI/ML Engineers: 120 positions open\n"
                "  - Cloud Solutions Architects: 85 positions open\n"
                "  - Security Specialists: 45 positions open\n"
                "  - Product Managers: 30 positions open\n\n"
                "Employee Retention Rate: 91.3%\n"
                "Average Tenure: 3.8 years\n"
                "Internal Promotion Rate: 28%",
                font_size=14)

    # ---- Slide 10: Risk Assessment ----
    slide10 = prs.slides.add_slide(blank_layout)
    add_title(slide10, "Risk Assessment Matrix")
    add_textbox(slide10, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5),
                "High Priority Risks:\n\n"
                "1. Semiconductor supply chain disruption (Medium-High)\n"
                "2. Regulatory changes in EU AI Act compliance (High)\n"
                "3. Talent competition from hyperscalers (Medium)\n"
                "4. Currency fluctuation impact on APAC revenue (Medium)",
                font_size=14)
    add_textbox(slide10, Inches(7), Inches(1.8), Inches(5.5), Inches(5),
                "Mitigation Strategies:\n\n"
                "1. Diversified supplier base across 3 continents\n"
                "2. Dedicated compliance team with quarterly audits\n"
                "3. Enhanced compensation and equity refresh program\n"
                "4. Natural hedging through regional pricing",
                font_size=14)

    # ---- Slide 11: Partnership Ecosystem ----
    slide11 = prs.slides.add_slide(blank_layout)
    add_title(slide11, "Strategic Partnership Ecosystem")
    add_textbox(slide11, Inches(0.8), Inches(1.8), Inches(11), Inches(5),
                "Tier 1 Partners (>$10M annual co-revenue):\n"
                "  Accenture, Deloitte, Wipro, Infosys\n\n"
                "Technology Alliance Partners:\n"
                "  NVIDIA, Intel, AMD, Arm Holdings\n\n"
                "Cloud Marketplace Presence:\n"
                "  AWS Marketplace (142 listings), Azure Marketplace (98 listings),\n"
                "  Google Cloud Marketplace (67 listings)\n\n"
                "Partner-sourced revenue: $87.4M (15% of total)",
                font_size=14)

    # ---- Slide 12: Sustainability ----
    slide12 = prs.slides.add_slide(blank_layout)
    add_title(slide12, "Sustainability & ESG Commitments")
    add_textbox(slide12, Inches(0.8), Inches(1.8), Inches(11), Inches(5),
                "Carbon Neutrality: On track for 2027 target\n"
                "  - Current emissions: 42,300 tCO2e (down 18% YoY)\n"
                "  - Renewable energy usage: 78% of total consumption\n\n"
                "Data Center Efficiency:\n"
                "  - PUE (Power Usage Effectiveness): 1.18 (industry avg: 1.58)\n"
                "  - Water recycling rate: 92%\n\n"
                "Social Impact:\n"
                "  - STEM education grants: $4.2M distributed\n"
                "  - Diversity in leadership: 38% (target: 45% by 2027)",
                font_size=14)

    # ---- Slide 13: Q4 Priorities ----
    slide13 = prs.slides.add_slide(blank_layout)
    add_title(slide13, "Q4 2025 Strategic Priorities")
    add_textbox(slide13, Inches(0.8), Inches(1.8), Inches(11), Inches(5),
                "1. Close 12 pending enterprise deals ($34M pipeline)\n"
                "2. Ship NexusAI 3.0 beta to design partners by Nov 15\n"
                "3. Complete SOC 2 Type II audit for SecureVault\n"
                "4. Finalize Series C investment in quantum computing startup\n"
                "5. Launch employee stock purchase plan (ESPP)\n"
                "6. Open Singapore regional headquarters\n"
                "7. Achieve ISO 27001 certification for APAC data centers",
                font_size=15)

    # ---- Slide 14: Thank You / Q&A ----
    slide14 = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(slide14, "Thank You")
    if 1 in slide14.placeholders:
        slide14.placeholders[1].text = (
            "Questions & Discussion\n\n"
            "Contact: strategy@meridiantech.com\n"
            "Next Review: February 2026"
        )
    else:
        add_textbox(slide14, Inches(2), Inches(4), Inches(9), Inches(2),
                    "Questions & Discussion\n\nContact: strategy@meridiantech.com\nNext Review: February 2026",
                    font_size=18, alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Number of slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
