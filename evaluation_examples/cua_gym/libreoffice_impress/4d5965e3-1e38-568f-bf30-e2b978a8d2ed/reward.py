"""
Reward Script: Add 'DRAFT' watermark across every slide
Task ID: impress_fix_071
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): All 14 slides have a 'DRAFT' text shape
  Component 2 (0.25): Font size is 72pt and color is gray #808080
  Component 3 (0.20): Watermark is rotated 45 degrees
  Component 4 (0.15): Semi-transparency (~50%) applied
  Component 5 (0.10): Watermark is behind other content (first in z-order)
"""

import os
from pptx import Presentation
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_071'
EXPECTED_SLIDE_COUNT = 14


def persist_app_state(domain: str):
    """Save any unsaved GUI edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def find_draft_shapes(slide):
    """Find all shapes with text 'DRAFT' on a slide."""
    draft_shapes = []
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip().upper() == 'DRAFT':
            draft_shapes.append(shape)
    return draft_shapes


def verify_task(file_path):
    """
    Verify DRAFT watermark is on every slide with correct properties.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    if num_slides == 0:
        print("FAIL: Presentation has no slides")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: All slides have a 'DRAFT' text shape (0.30 points)
    # This is the core requirement — DRAFT watermark must exist on every slide.
    try:
        slides_with_draft = 0
        for i, slide in enumerate(prs.slides):
            draft_shapes = find_draft_shapes(slide)
            if len(draft_shapes) > 0:
                slides_with_draft += 1
            else:
                print(f"FAIL: Slide {i+1} has no DRAFT watermark")

        # Award proportional credit based on fraction of slides covered
        if slides_with_draft == num_slides and num_slides >= EXPECTED_SLIDE_COUNT:
            print(f"PASS: Component 1 — All {num_slides} slides have DRAFT watermark (0.30 pts)")
            total_score += 0.30
        elif slides_with_draft > 0:
            fraction = slides_with_draft / num_slides
            partial = round(0.30 * fraction, 2)
            print(f"PARTIAL: Component 1 — {slides_with_draft}/{num_slides} slides have DRAFT watermark ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — No slides have DRAFT watermark (0 pts)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # For remaining components, collect properties from all DRAFT shapes
    all_draft_shapes = []
    for slide in prs.slides:
        drafts = find_draft_shapes(slide)
        all_draft_shapes.extend(drafts)

    if len(all_draft_shapes) == 0:
        # No DRAFT shapes found at all — remaining components all fail
        print("FAIL: Components 2-5 — No DRAFT shapes to evaluate")
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: Font size is 72pt (914400 EMU) and color is gray #808080 (0.25 points)
    try:
        correct_font_count = 0
        for shape in all_draft_shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip().upper() != 'DRAFT':
                        continue
                    size_ok = False
                    color_ok = False
                    # Check font size: 72pt = 914400 EMU
                    if run.font.size is not None:
                        # Allow small tolerance (within 1pt = 12700 EMU)
                        if abs(run.font.size - 914400) <= 12700:
                            size_ok = True
                    # Check color: #808080
                    try:
                        if run.font.color.type is not None:
                            rgb_str = str(run.font.color.rgb).upper()
                            if rgb_str == '808080':
                                color_ok = True
                    except Exception:
                        pass
                    if size_ok and color_ok:
                        correct_font_count += 1
                    else:
                        print(f"  font check: size_ok={size_ok} (size={run.font.size}), color_ok={color_ok}")

        if correct_font_count >= num_slides:
            print(f"PASS: Component 2 — All DRAFT shapes have 72pt gray #808080 font (0.25 pts)")
            total_score += 0.25
        elif correct_font_count > 0:
            fraction = correct_font_count / num_slides
            partial = round(0.25 * fraction, 2)
            print(f"PARTIAL: Component 2 — {correct_font_count}/{num_slides} DRAFT shapes have correct font ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No DRAFT shapes have correct 72pt gray font (0 pts)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Watermark is rotated 45 degrees (0.20 points)
    # Rotation is stored as rot attribute on the <p:sp> element in 60000ths of a degree
    # 45 degrees = 2700000
    try:
        rotated_count = 0
        for shape in all_draft_shapes:
            rot_str = shape._element.get('rot', '0')
            try:
                rot_val = int(rot_str)
            except (ValueError, TypeError):
                rot_val = 0
            # Allow tolerance: within 5 degrees (300000 units)
            if abs(rot_val - 2700000) <= 300000:
                rotated_count += 1
            else:
                print(f"  rotation check: rot_val={rot_val}, expected ~2700000 (45 deg)")

        if rotated_count >= num_slides:
            print(f"PASS: Component 3 — All DRAFT watermarks rotated ~45 degrees (0.20 pts)")
            total_score += 0.20
        elif rotated_count > 0:
            fraction = rotated_count / num_slides
            partial = round(0.20 * fraction, 2)
            print(f"PARTIAL: Component 3 — {rotated_count}/{num_slides} DRAFT shapes rotated correctly ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No DRAFT shapes are rotated 45 degrees (0 pts)")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Semi-transparency (~50%) applied (0.15 points)
    # Transparency is stored as <a:alpha val="50000"/> inside the font color solidFill
    try:
        transparent_count = 0
        for shape in all_draft_shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip().upper() != 'DRAFT':
                        continue
                    rPr = run._r.find(qn('a:rPr'))
                    if rPr is None:
                        continue
                    solidFill = rPr.find(qn('a:solidFill'))
                    if solidFill is None:
                        continue
                    # Check for alpha element in any color child
                    alpha_found = False
                    for color_elem in solidFill:
                        alpha_elem = color_elem.find(qn('a:alpha'))
                        if alpha_elem is not None:
                            alpha_val = alpha_elem.get('val', '100000')
                            try:
                                alpha_int = int(alpha_val)
                                # Semi-transparent means alpha < 100000 (fully opaque)
                                # Expect around 50000 (50%), allow 20000-80000 range
                                if 20000 <= alpha_int <= 80000:
                                    alpha_found = True
                                    break
                            except ValueError:
                                pass
                    if alpha_found:
                        transparent_count += 1
                    else:
                        print(f"  transparency check: no valid alpha found on a DRAFT shape")

        if transparent_count >= num_slides:
            print(f"PASS: Component 4 — All DRAFT watermarks have semi-transparency (0.15 pts)")
            total_score += 0.15
        elif transparent_count > 0:
            fraction = transparent_count / num_slides
            partial = round(0.15 * fraction, 2)
            print(f"PARTIAL: Component 4 — {transparent_count}/{num_slides} DRAFT shapes are semi-transparent ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — No DRAFT shapes have semi-transparency (0 pts)")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Watermark is behind other content (first in z-order) (0.10 points)
    # The DRAFT shape should be the first shape in the slide's shape tree (lowest z-order = behind)
    try:
        behind_count = 0
        for i, slide in enumerate(prs.slides):
            shapes_list = list(slide.shapes)
            if len(shapes_list) < 2:
                # Only one shape, cannot meaningfully check z-order
                draft_shapes = find_draft_shapes(slide)
                if draft_shapes:
                    behind_count += 1
                continue
            # Check if the first shape is the DRAFT watermark
            first_shape = shapes_list[0]
            if hasattr(first_shape, 'text') and first_shape.text.strip().upper() == 'DRAFT':
                behind_count += 1
            else:
                print(f"  z-order check: Slide {i+1} DRAFT is not the first shape")

        if behind_count >= num_slides:
            print(f"PASS: Component 5 — All DRAFT watermarks are behind content (0.10 pts)")
            total_score += 0.10
        elif behind_count > 0:
            fraction = behind_count / num_slides
            partial = round(0.10 * fraction, 2)
            print(f"PARTIAL: Component 5 — {behind_count}/{num_slides} DRAFT watermarks behind content ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 — No DRAFT watermarks are behind other content (0 pts)")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
