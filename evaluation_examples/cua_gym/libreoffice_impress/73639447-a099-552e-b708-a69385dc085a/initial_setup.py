"""
Initial Setup: Enable slide numbering starting from 0 in product_roadmap.pptx
Task ID: impress_slides_035
Domain: libreoffice_impress

Creates a 10-slide product roadmap presentation with slide numbering starting at 1 (default).
The task is to change it to start at 0.
"""

import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
TASK_ID = 'impress_slides_035'
DESKTOP_FILE = f'{DESKTOP}/product_roadmap.pptx'
INITIAL_FILE = f'{WORKDIR}/{TASK_ID}_initial.pptx'


def add_title_slide(prs, title, subtitle):
    layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, bullet_points):
    layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.text = bullet_points[0]
    for bp in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = bp
        p.level = 1
    return slide


def create_initial():
    os.makedirs(DESKTOP, exist_ok=True)

    prs = Presentation()
    # Standard widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Set firstSlideNum to 1 (default - the initial state before the task)
    prs._element.set('firstSlideNum', '1')

    # Slide 1: Title Slide
    add_title_slide(
        prs,
        "Nexaflow Technologies — Product Roadmap 2025",
        "Accelerating Innovation Across Every Touchpoint"
    )

    # Slide 2: Vision & Strategy
    add_content_slide(prs, "Vision & Strategic Pillars", [
        "Expand AI-assisted workflows to 3 core product lines",
        "Achieve 40% reduction in time-to-market for new features",
        "Deepen enterprise integrations with SAP, Salesforce, and Workday",
        "Grow developer ecosystem to 5,000 active API partners",
    ])

    # Slide 3: Q1 2025 Milestones
    add_content_slide(prs, "Q1 2025 — Foundation & Launch", [
        "Release NexaFlow Core v3.0 with real-time collaboration",
        "Launch Nexaflow Analytics Dashboard (public beta)",
        "Complete SOC 2 Type II audit certification",
        "Onboard first 10 enterprise pilot customers",
    ])

    # Slide 4: Q2 2025 Milestones
    add_content_slide(prs, "Q2 2025 — Scale & Integration", [
        "GA release of NexaFlow Mobile (iOS & Android)",
        "Integrate with Microsoft Teams and Slack (bi-directional sync)",
        "Ship AI-powered workflow suggestions engine",
        "Achieve 99.95% uptime SLA across all regions",
    ])

    # Slide 5: Q3 2025 Milestones
    add_content_slide(prs, "Q3 2025 — Intelligence & Automation", [
        "Deploy predictive analytics module for workflow bottlenecks",
        "Launch NexaFlow Marketplace with 50+ third-party connectors",
        "Introduce role-based access control v2 with attribute policies",
        "Expand data residency to EU, APAC, and LatAm regions",
    ])

    # Slide 6: Q4 2025 Milestones
    add_content_slide(prs, "Q4 2025 — Enterprise & Growth", [
        "Release NexaFlow Enterprise Suite with dedicated support tier",
        "Launch white-label partner program for system integrators",
        "Complete migration of legacy customers to v3 platform",
        "Publish public API v2 with GraphQL support",
    ])

    # Slide 7: Key Investments
    add_content_slide(prs, "Key Investment Areas — 2025", [
        "Engineering: +35 headcount across backend, ML, and platform teams",
        "Infrastructure: Migrate to multi-cloud Kubernetes on AWS + GCP",
        "Security: Zero-trust architecture rollout across all services",
        "Customer Success: Build dedicated enterprise onboarding program",
    ])

    # Slide 8: Risk & Mitigation
    add_content_slide(prs, "Risks & Mitigation Strategies", [
        "Risk: Talent acquisition delays — Mitigation: Partner with 3 recruiting firms",
        "Risk: Third-party API instability — Mitigation: Circuit-breaker patterns in v3",
        "Risk: Regulatory changes in EU — Mitigation: Dedicated compliance council",
        "Risk: Competitor feature parity — Mitigation: Accelerate unique AI differentiators",
    ])

    # Slide 9: Success Metrics
    add_content_slide(prs, "Success Metrics — End of 2025", [
        "ARR target: $48M (up from $31M in 2024)",
        "Net Revenue Retention: ≥ 118%",
        "Customer Satisfaction Score (CSAT): ≥ 87",
        "Feature delivery velocity: 2-week sprint cadence sustained",
    ])

    # Slide 10: Call to Action
    add_content_slide(prs, "Next Steps & Call to Action", [
        "Finalize Q1 sprint plans with all product leads by Jan 15",
        "Schedule quarterly business reviews with top 20 enterprise accounts",
        "Kick off 2025 engineering org design workshops",
        "Publish external product roadmap blog post by Feb 1",
    ])

    # Save to Desktop (the task file)
    prs.save(DESKTOP_FILE)
    print(f'Desktop file created: {DESKTOP_FILE}')

    # Also copy as initial tracking file
    shutil.copy(DESKTOP_FILE, INITIAL_FILE)
    print(f'Initial tracking file created: {INITIAL_FILE}')

    # Verify firstSlideNum
    from pptx import Presentation as P2
    prs_check = P2(DESKTOP_FILE)
    val = prs_check._element.get('firstSlideNum', '1')
    print(f'firstSlideNum in file: {val} (expected: 1)')
    print(f'Number of slides: {len(prs_check.slides)} (expected: 10)')


create_initial()
