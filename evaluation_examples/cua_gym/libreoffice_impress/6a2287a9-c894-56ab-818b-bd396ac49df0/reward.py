"""
Reward Script: Resize the image on slide 2 to fill the complete slide area, maintaining aspect ratio, centered.
Task ID: osworld_impress_image_fill_slide_008
Domain: libreoffice_impress
Scoring:
  Component 1: Image width equals slide width (0.4 pts)
  Component 2: Image height equals slide height (0.3 pts)
  Component 3: Image is centered — left and top at correct position (0.3 pts)
"""

import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_image_fill_slide_008'

# Tolerance for EMU comparisons (0.5%)
TOLERANCE = 0.005


def is_approx_equal(val1, val2, tolerance=TOLERANCE):
    """Check if two EMU values are approximately equal with relative tolerance."""
    if val1 == val2:
        return True
    if val1 == 0 or val2 == 0:
        return val1 == val2
    return abs(val1 - val2) / max(abs(val1), abs(val2)) <= tolerance


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Task: Resize the image on slide 2 to fill the complete slide area,
          maintaining aspect ratio, centered both vertically and horizontally.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"CRITICAL: Expected at least 2 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Get slide dimensions
    slide_width = prs.slide_width   # 9144000 EMU = 10 inches
    slide_height = prs.slide_height  # 6858000 EMU = 7.5 inches

    # Locate the picture on slide 2 (index 1)
    slide2 = prs.slides[1]
    picture_shape = None
    for shape in slide2.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture_shape = shape
            break

    if picture_shape is None:
        print("CRITICAL: No picture shape found on slide 2")
        print("REWARD: 0.0")
        return 0.0

    img_left = picture_shape.left
    img_top = picture_shape.top
    img_width = picture_shape.width
    img_height = picture_shape.height

    print(f"INFO: Slide dimensions: {slide_width} x {slide_height} EMU ({slide_width/914400:.4f} x {slide_height/914400:.4f} in)")
    print(f"INFO: Image position: left={img_left}, top={img_top}")
    print(f"INFO: Image size: width={img_width}, height={img_height} ({img_width/914400:.4f} x {img_height/914400:.4f} in)")

    # Component 1: Image width equals slide width (0.4 points)
    # Initial: width=4572000 (50% of slide) — FAILS
    # Golden:  width=9144000 (100% of slide) — PASSES
    try:
        width_ok = is_approx_equal(img_width, slide_width)
        if width_ok:
            print(f"PASS: Component 1 — Image width fills slide: {img_width} == {slide_width} EMU (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 — Image width {img_width} ({img_width/914400:.4f}in) != slide_width {slide_width} ({slide_width/914400:.4f}in)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Image height equals slide height (0.3 points)
    # Initial: height=3429000 (50% of slide) — FAILS
    # Golden:  height=6858000 (100% of slide) — PASSES
    try:
        height_ok = is_approx_equal(img_height, slide_height)
        if height_ok:
            print(f"PASS: Component 2 — Image height fills slide: {img_height} == {slide_height} EMU (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 — Image height {img_height} ({img_height/914400:.4f}in) != slide_height {slide_height} ({slide_height/914400:.4f}in)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Image is centered — left and top at centered position (0.3 points)
    # expected_left = (slide_width - img_width) // 2
    # expected_top  = (slide_height - img_height) // 2
    # Initial: left=457200, expected=2286000 — FAILS
    # Golden:  left=0, expected=0 — PASSES
    try:
        expected_left = (slide_width - img_width) // 2
        expected_top = (slide_height - img_height) // 2
        left_ok = (img_left == 0) if expected_left == 0 else is_approx_equal(img_left, expected_left)
        top_ok = (img_top == 0) if expected_top == 0 else is_approx_equal(img_top, expected_top)
        if left_ok and top_ok:
            print(f"PASS: Component 3 — Image centered: left={img_left} (exp {expected_left}), top={img_top} (exp {expected_top}) (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 3 — Not centered: left={img_left} (exp {expected_left}), top={img_top} (exp {expected_top})")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
