"""
Initial Setup: 4-slide event promotion deck with image on slide 2 at ~50% size
Task ID: osworld_impress_image_fill_slide_008
Domain: libreoffice_impress
"""

import os
import io
import shlex
import subprocess
import time

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_image_fill_slide_008'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
IMG_PATH = f'{WORKDIR}/{TASK_ID}_event_photo.png'

def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_event_photo():
    """Create a realistic-looking event photo image using PIL."""
    from PIL import Image, ImageDraw, ImageFont
    import math

    # Create a 800x600 image (4:3 aspect ratio for the event photo)
    width, height = 800, 600
    img = Image.new('RGB', (width, height), color=(30, 50, 90))
    draw = ImageDraw.Draw(img)

    # Draw a gradient sky background
    for y in range(height // 2):
        ratio = y / (height // 2)
        r = int(30 + ratio * 60)
        g = int(50 + ratio * 80)
        b = int(90 + ratio * 120)
        draw.line([(0, y), (width, y)], fill=(r, g, b))

    # Draw ground/stage area
    for y in range(height // 2, height):
        ratio = (y - height // 2) / (height // 2)
        r = int(80 + ratio * 40)
        g = int(60 + ratio * 30)
        b = int(20 + ratio * 10)
        draw.line([(0, y), (width, y)], fill=(r, g, b))

    # Draw stage lights
    light_colors = [(255, 200, 50), (50, 200, 255), (255, 80, 80), (80, 255, 120)]
    light_positions = [(150, 80), (300, 60), (500, 65), (650, 75)]
    for (lx, ly), lc in zip(light_positions, light_colors):
        # Light beam
        for radius in range(60, 0, -5):
            alpha = int(30 + (60 - radius) * 2)
            draw.ellipse(
                [(lx - radius, ly - radius // 2), (lx + radius, ly + radius // 2)],
                fill=(*lc, min(255, alpha))
            )
        # Light source
        draw.ellipse([(lx - 15, ly - 15), (lx + 15, ly + 15)], fill=(255, 255, 200))

    # Draw stage platform
    draw.rectangle([(50, height // 2 + 60), (width - 50, height - 30)], fill=(60, 40, 20))
    draw.rectangle([(50, height // 2 + 60), (width - 50, height // 2 + 75)], fill=(90, 70, 40))

    # Draw crowd silhouettes
    crowd_y = height // 2 + 20
    for x in range(20, width - 20, 25):
        person_height = 30 + (x % 15)
        head_y = crowd_y - person_height
        draw.ellipse([(x - 8, head_y - 10), (x + 8, head_y + 5)], fill=(20, 20, 20))
        draw.rectangle([(x - 6, head_y + 5), (x + 6, crowd_y)], fill=(20, 20, 20))

    # Add event banner text
    draw.rectangle([(100, 150), (700, 230)], fill=(180, 20, 40))
    draw.rectangle([(102, 152), (698, 228)], fill=(200, 30, 50))

    # Draw banner text using basic shapes (since no font file guaranteed)
    banner_text = "ANNUAL TECH SUMMIT 2025"
    # Use a simple block letter approach with rectangle
    text_x, text_y = 130, 168
    # Draw each "pixel" of text as small blocks for title
    draw.text_if_possible = False

    # Try to add text with default font
    try:
        draw.text((text_x, text_y), banner_text, fill=(255, 255, 255))
        draw.text((110, 195), "Innovation | Inspiration | Impact", fill=(255, 220, 180))
    except Exception:
        pass

    # Add some decorative elements
    for i in range(10):
        star_x = 60 + i * 70
        star_y = 120 + (i % 3) * 15
        draw.ellipse([(star_x - 3, star_y - 3), (star_x + 3, star_y + 3)], fill=(255, 255, 100))

    # Add venue text at bottom
    draw.rectangle([(0, height - 50), (width, height)], fill=(10, 10, 30))
    try:
        draw.text((20, height - 35), "Grand Convention Center  |  March 15-17, 2025  |  San Francisco, CA", fill=(200, 200, 200))
    except Exception:
        pass

    return img


def create_initial():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import io

    # Create event photo image and save to temp file
    photo_img = create_event_photo()
    photo_img.save(IMG_PATH)
    print(f'Event photo created: {IMG_PATH}')

    prs = Presentation()
    # Use widescreen 10" x 7.5" (python-pptx default)
    slide_w = prs.slide_width   # 9144000 EMU = 10 inches
    slide_h = prs.slide_height  # 6858000 EMU = 7.5 inches

    # ---- SLIDE 1: Title slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
    slide1.shapes.title.text = "Annual Tech Summit 2025"
    slide1.placeholders[1].text = "Innovation | Inspiration | Impact\nMarch 15-17, 2025 | San Francisco, CA"

    # Style title
    title_frame = slide1.shapes.title.text_frame
    for para in title_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1A, 0x3A, 0x6B)

    # Style subtitle
    subtitle_frame = slide1.placeholders[1].text_frame
    for para in subtitle_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # Background for slide 1
    fill1 = slide1.background.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFF)

    # ---- SLIDE 2: Event photo slide (image at ~50% of slide size) ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title text box at top
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Event Highlights"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x3A, 0x6B)

    # Add image at ~50% of slide size, positioned upper-left area (NOT centered, NOT filling)
    # Slide: 10" x 7.5" | 50% => 5" x 3.75"
    img_w = Inches(5.0)   # 50% of slide width
    img_h = Inches(3.75)  # 50% of slide height (keeping 4:3 ratio from original 800x600)
    img_left = Inches(0.5)  # Positioned at left, not centered
    img_top = Inches(1.2)   # Below title
    pic = slide2.shapes.add_picture(IMG_PATH, img_left, img_top, img_w, img_h)

    # Add caption below image
    caption_box = slide2.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.5))
    cf = caption_box.text_frame
    cp = cf.paragraphs[0]
    cp.text = "Thousands of attendees gathered for keynote sessions and workshops"
    cp_run = cp.runs[0]
    cp_run.font.size = Pt(14)
    cp_run.font.italic = True
    cp_run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # Background for slide 2
    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # ---- SLIDE 3: Event details / schedule ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content layout
    slide3.shapes.title.text = "Conference Schedule"

    # Style slide 3 title
    s3_title = slide3.shapes.title.text_frame
    for para in s3_title.paragraphs:
        for run in para.runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1A, 0x3A, 0x6B)

    # Add schedule content
    content_ph = slide3.placeholders[1]
    tf3 = content_ph.text_frame
    tf3.text = "Day 1 – Keynote & Opening Ceremonies"

    schedule_items = [
        "Day 2 – Workshops & Technical Sessions",
        "Day 3 – Panel Discussions & Networking",
        "Over 120 speakers from 40+ countries",
        "50+ hands-on workshops and labs",
        "Startup Expo with 200+ exhibitors",
    ]
    for item in schedule_items:
        p_new = tf3.add_paragraph()
        p_new.text = item
        p_new.level = 1
        for run in p_new.runs:
            run.font.size = Pt(18)

    # Background for slide 3
    fill3 = slide3.background.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xFF)

    # ---- SLIDE 4: Closing / Call to action ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
    slide4.shapes.title.text = "Join Us at the Summit!"
    slide4.placeholders[1].text = (
        "Register now at www.techsummit2025.com\n"
        "Early bird pricing ends February 28\n"
        "Contact: info@techsummit2025.com"
    )

    # Style slide 4
    s4_title = slide4.shapes.title.text_frame
    for para in s4_title.paragraphs:
        for run in para.runs:
            run.font.size = Pt(36)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    s4_sub = slide4.placeholders[1].text_frame
    for para in s4_sub.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0xE0, 0xE8, 0xFF)

    # Dark background for closing slide
    fill4 = slide4.background.fill
    fill4.solid()
    fill4.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x6B)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup — open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
