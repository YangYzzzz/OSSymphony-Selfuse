"""
Initial Setup: Create a standard 4:3 presentation with 8 slides for board meeting
Task ID: impress_gf5_010
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_010'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=None):
    """Helper to add a textbox with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()

    # Standard 4:3 dimensions: 25.4cm x 19.05cm (10 x 7.5 inches)
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(19.05)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x00, 0x33, 0x66)
    add_textbox(slide1, Inches(1), Inches(1.5), Inches(8), Inches(1.5),
                "Meridian Technologies Inc.", font_size=36, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(1), Inches(3.2), Inches(8), Inches(1),
                "Q4 2025 Board of Directors Review", font_size=24,
                color=RGBColor(0xCC, 0xDD, 0xEE), alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(1), Inches(5.0), Inches(8), Inches(0.6),
                "December 15, 2025  |  Presented by Elena Vasquez, CEO",
                font_size=14, color=RGBColor(0x99, 0xBB, 0xCC),
                alignment=PP_ALIGN.CENTER)

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide2, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                "Agenda", font_size=32, bold=True,
                color=RGBColor(0x00, 0x33, 0x66))
    agenda_items = [
        "1.  Financial Performance Overview",
        "2.  Revenue Breakdown by Segment",
        "3.  Key Product Milestones",
        "4.  Customer Acquisition & Retention",
        "5.  Operational Highlights",
        "6.  Strategic Outlook for 2026",
        "7.  Q&A / Open Discussion",
    ]
    txBox = slide2.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(agenda_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.space_after = Pt(12)
        run = p.runs[0]
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 3: Financial Performance ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide3, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                "Financial Performance Overview", font_size=28, bold=True,
                color=RGBColor(0x00, 0x33, 0x66))
    # Key metrics
    metrics = [
        ("Total Revenue", "$142.8M", "+18.3% YoY"),
        ("Gross Margin", "67.2%", "+2.1 pp"),
        ("EBITDA", "$38.5M", "+22.7% YoY"),
        ("Free Cash Flow", "$24.1M", "+15.9% YoY"),
    ]
    for i, (label, value, change) in enumerate(metrics):
        y = Inches(1.4) + Inches(1.2) * i
        add_textbox(slide3, Inches(0.8), y, Inches(3), Inches(0.5),
                    label, font_size=16, color=RGBColor(0x66, 0x66, 0x66))
        add_textbox(slide3, Inches(4.0), y, Inches(2.5), Inches(0.5),
                    value, font_size=20, bold=True,
                    color=RGBColor(0x00, 0x33, 0x66))
        add_textbox(slide3, Inches(6.5), y, Inches(2.5), Inches(0.5),
                    change, font_size=16,
                    color=RGBColor(0x22, 0x8B, 0x22))

    # --- Slide 4: Revenue Breakdown ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide4, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                "Revenue Breakdown by Segment", font_size=28, bold=True,
                color=RGBColor(0x00, 0x33, 0x66))
    table_shape = slide4.shapes.add_table(6, 4, Inches(0.5), Inches(1.3),
                                           Inches(9), Inches(4.5))
    table = table_shape.table
    headers = ["Segment", "Q4 Revenue", "% of Total", "YoY Growth"]
    data_rows = [
        ["Cloud Services", "$58.4M", "40.9%", "+28.1%"],
        ["Enterprise Software", "$42.1M", "29.5%", "+12.4%"],
        ["Professional Services", "$24.7M", "17.3%", "+8.7%"],
        ["Hardware Solutions", "$11.2M", "7.8%", "+3.2%"],
        ["Licensing & Other", "$6.4M", "4.5%", "+5.1%"],
    ]
    for col, h in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x00, 0x33, 0x66)
    for r, row_data in enumerate(data_rows, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 5: Product Milestones ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide5, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                "Key Product Milestones", font_size=28, bold=True,
                color=RGBColor(0x00, 0x33, 0x66))
    milestones = [
        "Meridian Cloud Platform 3.0 launched (October) - 15,000+ enterprise signups in first month",
        "AI-Powered Analytics Suite reached GA - adopted by 340 Fortune 500 clients",
        "Security Compliance Module achieved FedRAMP High authorization",
        "Mobile SDK 2.0 released with 47% performance improvement",
        "Strategic partnership with Hyperion Systems for integrated ERP workflows",
    ]
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, m in enumerate(milestones):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"\u2022  {m}"
        p.space_after = Pt(14)
        run = p.runs[0]
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 6: Customer Metrics ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide6, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                "Customer Acquisition & Retention", font_size=28, bold=True,
                color=RGBColor(0x00, 0x33, 0x66))
    cust_metrics = [
        ("New Enterprise Clients", "127", "vs. 98 in Q3"),
        ("Net Revenue Retention", "118%", "up from 114%"),
        ("Customer Satisfaction (NPS)", "72", "Industry avg: 54"),
        ("Annual Churn Rate", "4.2%", "down from 5.8%"),
    ]
    for i, (label, value, note) in enumerate(cust_metrics):
        y = Inches(1.4) + Inches(1.2) * i
        add_textbox(slide6, Inches(0.8), y, Inches(3.5), Inches(0.5),
                    label, font_size=16, color=RGBColor(0x66, 0x66, 0x66))
        add_textbox(slide6, Inches(4.5), y, Inches(1.5), Inches(0.5),
                    value, font_size=22, bold=True,
                    color=RGBColor(0x00, 0x33, 0x66))
        add_textbox(slide6, Inches(6.2), y, Inches(3), Inches(0.5),
                    note, font_size=14,
                    color=RGBColor(0x88, 0x88, 0x88))

    # --- Slide 7: Operational Highlights ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide7, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                "Operational Highlights", font_size=28, bold=True,
                color=RGBColor(0x00, 0x33, 0x66))
    ops = [
        "Headcount grew to 2,847 employees (+312 in Q4), with key hires in AI/ML division",
        "New Singapore data center operational - reduces APAC latency by 40%",
        "SOC 2 Type II certification renewed with zero findings",
        "Infrastructure costs reduced 11% through container optimization initiative",
        "Employee engagement score: 4.3/5.0 (up from 4.1 in Q3)",
    ]
    txBox = slide7.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(ops):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"\u2022  {item}"
        p.space_after = Pt(14)
        run = p.runs[0]
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 8: Strategic Outlook ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide8.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x00, 0x33, 0x66)
    add_textbox(slide8, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                "Strategic Outlook 2026", font_size=32, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF))
    outlook = [
        "Target $175M revenue (+22.5% growth)",
        "Launch Meridian AI Copilot for enterprise workflows",
        "Expand into European market with Frankfurt data center",
        "Pursue strategic acquisition in cybersecurity space ($30-50M range)",
        "Achieve carbon-neutral operations by Q3 2026",
    ]
    txBox = slide8.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(outlook):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"\u2022  {item}"
        p.space_after = Pt(16)
        run = p.runs[0]
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
