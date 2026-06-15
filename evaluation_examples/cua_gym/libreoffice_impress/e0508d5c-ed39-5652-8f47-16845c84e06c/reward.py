"""
Reward Script: Change slide size from 4:3 to 16:9 widescreen
Task ID: impress_gf5_010
Domain: libreoffice_impress
Scoring:
  Precondition gate: Height is 19.05cm AND 8 slides present (no points, early exit if broken)
  Component 1 (0.6): Slide width is widescreen 33.87cm (12192000 EMU +/- tolerance)
  Component 2 (0.4): Content was scaled to fit widescreen (shape widths increased proportionally)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_010'

# Expected EMU values
WIDESCREEN_WIDTH_EMU = 12192000   # 33.87cm ~ 13.33 inches (16:9)
STANDARD_WIDTH_EMU = 9144000      # 25.40cm ~ 10 inches (4:3)
EXPECTED_HEIGHT_EMU = 6858000     # 19.05cm ~ 7.5 inches
EXPECTED_SLIDE_COUNT = 8

# Tolerance: 0.5% relative
def approx_eq(actual, expected, tol=0.005):
    if expected == 0:
        return actual == 0
    return abs(actual - expected) / abs(expected) <= tol


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition gate: height must remain 19.05cm and all 8 slides must be present
    # These are pre-existing properties, so they serve as gates, not scoring components.
    slide_height = prs.slide_height
    slide_count = len(prs.slides)
    if not approx_eq(slide_height, EXPECTED_HEIGHT_EMU):
        print(f"GATE FAIL: Slide height changed to {slide_height} EMU ({slide_height/360000:.2f} cm), expected {EXPECTED_HEIGHT_EMU} EMU (19.05cm)")
        print("REWARD: 0.0")
        return 0.0
    if slide_count != EXPECTED_SLIDE_COUNT:
        print(f"GATE FAIL: Slide count is {slide_count}, expected {EXPECTED_SLIDE_COUNT}")
        print("REWARD: 0.0")
        return 0.0
    print(f"GATE PASS: Height = {slide_height/360000:.2f} cm, Slides = {slide_count}")

    # Component 1: Slide width is widescreen ~33.87cm / 12192000 EMU (0.6 points)
    # This FAILS on initial (25.40cm) and PASSES on golden (33.87cm)
    try:
        slide_width = prs.slide_width
        if approx_eq(slide_width, WIDESCREEN_WIDTH_EMU):
            print(f"PASS: Component 1 — Slide width is widescreen ({slide_width} EMU ~ {slide_width/360000:.2f} cm) (0.6 pts)")
            total_score += 0.6
        else:
            print(f"FAIL: Component 1 — Expected width ~{WIDESCREEN_WIDTH_EMU} EMU (33.87cm), found {slide_width} EMU ({slide_width/360000:.2f} cm)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Content was scaled to fit widescreen (0.4 points)
    # The "Scale to new format" option scales shape positions and widths proportionally.
    # Initial slide 1 shapes had width=7315200 EMU. Scaled by 12192000/9144000 ~ 1.333 -> ~9753600.
    # This FAILS on initial (shapes at original width) and PASSES on golden (shapes scaled wider).
    try:
        INITIAL_SLIDE1_SHAPE_WIDTH = 7315200  # original 4:3 width of TextBox 1 on slide 1
        SCALE_FACTOR = WIDESCREEN_WIDTH_EMU / STANDARD_WIDTH_EMU  # ~1.3333

        slide1 = prs.slides[0]
        widest_shape = max(s.width for s in slide1.shapes)

        # If content was scaled, widest shape should be ~9753600 EMU (7315200 * 1.333)
        expected_scaled_width = int(INITIAL_SLIDE1_SHAPE_WIDTH * SCALE_FACTOR)

        # Check: the widest shape width should be notably larger than the original 4:3 width
        # and approximately match the scaled width (within 5% tolerance for rounding)
        if widest_shape > INITIAL_SLIDE1_SHAPE_WIDTH * 1.15 and approx_eq(widest_shape, expected_scaled_width, tol=0.05):
            print(f"PASS: Component 2 — Content scaled: widest shape on slide 1 = {widest_shape} EMU (expected ~{expected_scaled_width}) (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 2 — Content not scaled properly: widest shape on slide 1 = {widest_shape} EMU, expected ~{expected_scaled_width} EMU")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
persist_app_state()

if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
