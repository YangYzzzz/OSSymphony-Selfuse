"""
Initial Setup: Create a 12-slide technical presentation deck
Task ID: impress_tm_068
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_068'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=None):
    """Helper to add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16):
    """Add a bullet list text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(font_size)
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0D, 0x25, 0x3F)
    add_textbox(slide1, Inches(1.5), Inches(1.5), Inches(10), Inches(2),
                "Distributed Systems Architecture", font_size=36, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(2), Inches(3.8), Inches(9), Inches(1),
                "A Comprehensive Technical Overview for Engineering Teams",
                font_size=20, color=RGBColor(0xB0, 0xC4, 0xDE),
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(3), Inches(5.5), Inches(7), Inches(0.8),
                "Dr. Elena Vasquez  |  Principal Architect  |  March 2026",
                font_size=14, color=RGBColor(0x80, 0x99, 0xB3),
                alignment=PP_ALIGN.CENTER)

    # ---- Slide 2: Agenda ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide2, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
                "Agenda", font_size=32, bold=True,
                color=RGBColor(0x0D, 0x25, 0x3F))
    add_bullet_list(slide2, Inches(1), Inches(1.6), Inches(10), Inches(5), [
        "1. System Architecture Overview",
        "2. Microservices Design Patterns",
        "3. Data Consistency Strategies",
        "4. Event-Driven Communication",
        "5. Observability & Monitoring",
        "6. Security Considerations",
        "7. Performance Benchmarks",
        "8. Deployment Pipeline",
        "9. Disaster Recovery Planning",
        "10. Q&A and Next Steps",
    ], font_size=18)

    # ---- Slide 3: Team Introduction ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide3, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
                "Core Engineering Team", font_size=32, bold=True,
                color=RGBColor(0x0D, 0x25, 0x3F))
    team_members = [
        "Elena Vasquez - Principal Architect (12 years distributed systems)",
        "James Nakamura - Senior Backend Engineer (Kafka, gRPC specialist)",
        "Priya Sharma - DevOps Lead (Kubernetes, Terraform, CI/CD)",
        "Carlos Mendez - Security Engineer (Zero-trust architecture)",
        "Wei Lin - Performance Engineer (Load testing, profiling)",
    ]
    add_bullet_list(slide3, Inches(1), Inches(1.6), Inches(11), Inches(5),
                    team_members, font_size=16)

    # ---- Slide 4: System Architecture Overview ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide4, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
                "System Architecture Overview", font_size=32, bold=True,
                color=RGBColor(0x0D, 0x25, 0x3F))
    add_bullet_list(slide4, Inches(1), Inches(1.6), Inches(5.5), Inches(5), [
        "3-tier architecture with API Gateway",
        "12 microservices across 4 bounded contexts",
        "PostgreSQL primary + Redis caching layer",
        "RabbitMQ for async message brokering",
        "Consul for service discovery & config",
    ], font_size=16)
    add_textbox(slide4, Inches(7), Inches(2), Inches(5.5), Inches(4),
                "Peak throughput: 45,000 req/s\nP99 latency: 23ms\n"
                "Uptime SLA: 99.97%\nData centers: us-east-1, eu-west-1",
                font_size=14, color=RGBColor(0x33, 0x33, 0x33))

    # ---- Slide 5: Microservices Design Patterns ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide5, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
                "Microservices Design Patterns", font_size=32, bold=True,
                color=RGBColor(0x0D, 0x25, 0x3F))
    add_bullet_list(slide5, Inches(1), Inches(1.6), Inches(11), Inches(5), [
        "Circuit Breaker: Hystrix-based with 5s timeout, 60% threshold",
        "Saga Pattern: Choreography for order processing pipeline",
        "CQRS: Separate read/write models for inventory service",
        "Sidecar Proxy: Envoy for mTLS and traffic shaping",
        "Bulkhead: Thread pool isolation per downstream dependency",
        "Retry with exponential backoff: max 3 retries, 200ms base",
    ], font_size=16)

    # ---- Slide 6: Data Consistency Strategies ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide6, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
                "Data Consistency Strategies", font_size=32, bold=True,
                color=RGBColor(0x0D, 0x25, 0x3F))
    add_bullet_list(slide6, Inches(1), Inches(1.6), Inches(5.5), Inches(5), [
        "Eventual consistency via event sourcing",
        "Two-phase commit for financial transactions",
        "Conflict-free replicated data types (CRDTs)",
        "Change Data Capture with Debezium",
        "Idempotency keys for all write operations",
    ], font_size=16)
    add_textbox(slide6, Inches(7), Inches(2), Inches(5.5), Inches(4),
                "Consistency model by service:\n"
                "  Payment: Strong (2PC)\n"
                "  Inventory: Eventual (CDC)\n"
                "  User Profile: Eventual (CRDT)\n"
                "  Order: Saga-based",
                font_size=14, color=RGBColor(0x33, 0x33, 0x33))

    # ---- Slide 7: Event-Driven Communication ----
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide7, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
                "Event-Driven Communication", font_size=32, bold=True,
                color=RGBColor(0x0D, 0x25, 0x3F))
    add_bullet_list(slide7, Inches(1), Inches(1.6), Inches(11), Inches(5), [
        "Apache Kafka: 6 brokers, 3 data centers, 7-day retention",
        "Schema Registry with Avro for contract evolution",
        "Dead letter queues for poison message handling",
        "Event replay capability for rebuilding read models",
        "Consumer group lag alerting (threshold: 10,000 messages)",
        "Exactly-once semantics via transactional producers",
    ], font_size=16)

    # ---- Slide 8: Observability & Monitoring ----
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide8, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
                "Observability & Monitoring", font_size=32, bold=True,
                color=RGBColor(0x0D, 0x25, 0x3F))
    add_bullet_list(slide8, Inches(1), Inches(1.6), Inches(5.5), Inches(5), [
        "Distributed tracing: Jaeger with OpenTelemetry SDK",
        "Metrics: Prometheus + Grafana (450+ custom dashboards)",
        "Logging: ELK stack with structured JSON logging",
        "Alerting: PagerDuty integration, 4 severity tiers",
        "SLO tracking: 99.95% availability, <50ms P95",
    ], font_size=16)
    add_textbox(slide8, Inches(7), Inches(2), Inches(5.5), Inches(4),
                "Monthly metrics volume:\n"
                "  Traces: 2.3 billion spans\n"
                "  Metrics: 890M time series\n"
                "  Logs: 14 TB compressed\n"
                "  Alerts triggered: ~340/month",
                font_size=14, color=RGBColor(0x33, 0x33, 0x33))

    # ---- Slide 9: Security Considerations ----
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide9, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
                "Security Considerations", font_size=32, bold=True,
                color=RGBColor(0x0D, 0x25, 0x3F))
    add_bullet_list(slide9, Inches(1), Inches(1.6), Inches(11), Inches(5), [
        "Zero-trust networking: All inter-service communication via mTLS",
        "OAuth 2.0 + OIDC for external API authentication",
        "Vault-managed secrets rotation every 24 hours",
        "Network policies: default-deny with explicit allowlists",
        "Container image scanning in CI pipeline (Trivy + Snyk)",
        "Quarterly penetration testing by external Red Team",
    ], font_size=16)

    # ---- Slide 10: Performance Benchmarks ----
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide10, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
                "Performance Benchmarks", font_size=32, bold=True,
                color=RGBColor(0x0D, 0x25, 0x3F))
    # Table with benchmark data
    table_shape = slide10.shapes.add_table(6, 4, Inches(1), Inches(1.8),
                                           Inches(11), Inches(4))
    table = table_shape.table
    headers = ["Service", "P50 Latency", "P99 Latency", "Throughput (req/s)"]
    data = [
        ["API Gateway", "4ms", "18ms", "45,000"],
        ["Auth Service", "8ms", "35ms", "12,000"],
        ["Order Service", "12ms", "67ms", "8,500"],
        ["Payment Service", "45ms", "180ms", "3,200"],
        ["Search Service", "22ms", "95ms", "15,000"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(13)

    # ---- Slide 11: Deployment Pipeline ----
    slide11 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide11, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
                "Deployment Pipeline", font_size=32, bold=True,
                color=RGBColor(0x0D, 0x25, 0x3F))
    add_bullet_list(slide11, Inches(1), Inches(1.6), Inches(11), Inches(5), [
        "Git push -> GitHub Actions CI (lint, test, build): ~4 min",
        "Container build + push to ECR: ~2 min",
        "ArgoCD sync to staging cluster: ~1 min",
        "Automated integration tests in staging: ~8 min",
        "Canary deployment to production (5% -> 25% -> 100%): ~30 min",
        "Rollback automation: <60s detection, <90s full rollback",
    ], font_size=16)

    # ---- Slide 12: Q&A / Next Steps ----
    slide12 = prs.slides.add_slide(prs.slide_layouts[5])
    fill = slide12.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0D, 0x25, 0x3F)
    add_textbox(slide12, Inches(2), Inches(2), Inches(9), Inches(2),
                "Questions & Next Steps", font_size=36, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_textbox(slide12, Inches(2), Inches(4.2), Inches(9), Inches(2),
                "Contact: elena.vasquez@techcorp.io\n"
                "Confluence: /wiki/distributed-systems-arch\n"
                "Slack: #arch-distributed-systems",
                font_size=18, color=RGBColor(0xB0, 0xC4, 0xDE),
                alignment=PP_ALIGN.CENTER)

    # No custom slideshows - this is the initial state
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
