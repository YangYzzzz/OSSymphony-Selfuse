"""
Reward Script: Apply strikethrough to 'Old Price: $49.99' and make 'New Price: $29.99' bold+green on slide 3
Task ID: impress_tct_084
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Strikethrough on 'Old Price: $49.99'
  Component 2 (0.3): Bold on 'New Price: $29.99'
  Component 3 (0.3): Green (#2E7D32) color on 'New Price: $29.99'
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_084'


def persist_app_state(domain: str):
    """Save any unsaved GUI edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
    except ImportError:
        print("CRITICAL: python-pptx not installed")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"CRITICAL: Expected at least 3 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[2]  # Slide 3 (0-indexed)

    # Find the text box containing the price lines
    old_price_run = None
    new_price_run = None

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                text = (run.text or "").strip()
                if "Old Price" in text and "$49.99" in text:
                    old_price_run = run
                if "New Price" in text and "$29.99" in text:
                    new_price_run = run

    if old_price_run is None:
        print("CRITICAL: Could not find 'Old Price: $49.99' text on slide 3")
        print("REWARD: 0.0")
        return 0.0

    if new_price_run is None:
        print("CRITICAL: Could not find 'New Price: $29.99' text on slide 3")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Strikethrough on 'Old Price: $49.99' (0.4 points)
    try:
        strike_val = old_price_run.font._element.attrib.get('strike', 'noStrike')
        if strike_val == 'sngStrike':
            print(f"PASS: Component 1 - 'Old Price' has strikethrough (strike={strike_val}) (0.4 pts)")
            total_score += 0.4
        elif strike_val == 'dblStrike':
            # Double strikethrough is close but not exact
            print(f"PARTIAL: Component 1 - 'Old Price' has double strikethrough instead of single (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 1 - 'Old Price' strike={strike_val}, expected 'sngStrike'")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Bold on 'New Price: $29.99' (0.3 points)
    try:
        bold_val = new_price_run.font.bold
        # Normalize None to False
        actual_bold = False if bold_val is None else bold_val
        if actual_bold is True:
            print(f"PASS: Component 2 - 'New Price' is bold (bold={bold_val}) (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 - 'New Price' bold={bold_val}, expected True")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Green color (#2E7D32) on 'New Price: $29.99' (0.3 points)
    try:
        color_type = new_price_run.font.color.type
        if color_type is not None:
            actual_rgb = str(new_price_run.font.color.rgb).upper()
            expected_rgb = "2E7D32"
            if actual_rgb == expected_rgb:
                print(f"PASS: Component 3 - 'New Price' color is #{actual_rgb} (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 3 - 'New Price' color is #{actual_rgb}, expected #{expected_rgb}")
        else:
            print(f"FAIL: Component 3 - 'New Price' has no explicit color set (inherited/theme)")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
