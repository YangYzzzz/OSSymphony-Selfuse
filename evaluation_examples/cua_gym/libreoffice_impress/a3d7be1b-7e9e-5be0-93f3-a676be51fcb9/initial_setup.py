"""
Initial Setup: Apply strikethrough and bold/green formatting on slide 3
Task ID: impress_tct_084
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_084'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Summer Promotion Campaign"
    slide1.placeholders[1].text = "Exclusive Deals for Our Valued Customers"

    # --- Slide 2: Campaign Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Campaign Overview"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Duration: June 15 - July 31, 2025"
    p2a = tf2.add_paragraph()
    p2a.text = "Target Audience: Premium subscribers and returning customers"
    p2b = tf2.add_paragraph()
    p2b.text = "Channels: Email marketing, social media, in-store signage"
    p2c = tf2.add_paragraph()
    p2c.text = "Expected Reach: 250,000+ potential customers"

    # --- Slide 3: Pricing (the task-relevant slide) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Title text box
    title_box = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "Special Pricing"
    title_p.alignment = PP_ALIGN.LEFT
    title_run = title_p.runs[0]
    title_run.font.size = Pt(28)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Price text box with two lines — both regular black 16pt (no strikethrough, no bold, no green)
    price_box = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(2))
    price_tf = price_box.text_frame
    price_tf.word_wrap = True

    # Line 1: Old Price
    p_old = price_tf.paragraphs[0]
    p_old.alignment = PP_ALIGN.LEFT
    run_old = p_old.add_run()
    run_old.text = "Old Price: $49.99"
    run_old.font.size = Pt(16)
    run_old.font.bold = False
    run_old.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # Line 2: New Price
    p_new = price_tf.add_paragraph()
    p_new.alignment = PP_ALIGN.LEFT
    run_new = p_new.add_run()
    run_new.text = "New Price: $29.99"
    run_new.font.size = Pt(16)
    run_new.font.bold = False
    run_new.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # --- Slide 4: Contact ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Get In Touch"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Email: promotions@summerdeal.com"
    p4a = tf4.add_paragraph()
    p4a.text = "Phone: +1 (555) 234-8900"
    p4b = tf4.add_paragraph()
    p4b.text = "Website: www.summerdeal.com/promo"
    p4c = tf4.add_paragraph()
    p4c.text = "Follow us @SummerDeals on all platforms"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
