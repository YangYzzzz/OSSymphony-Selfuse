"""
Reward Script: Template Library Presentation
Task ID: impress_wf_039
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): Slide count >= 12
  Component 2 (0.20): Title Slide layout pattern (slides with centered bold title + centered subtitle)
  Component 3 (0.15): Content layout pattern (slides with title bar rectangle + body text)
  Component 4 (0.15): Two Column layout pattern (slides with title + two side-by-side text areas)
  Component 5 (0.10): Image Focus layout pattern (slides with large placeholder + caption)
  Component 6 (0.10): Quote layout pattern (slides with italic centered text + attribution)
  Component 7 (0.10): Section Divider layout pattern (solid background + centered white text)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_039'


def get_text_shapes(slide):
    """Get all shapes that have text frames."""
    results = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            results.append(shape)
    return results


def get_nonempty_text(shape):
    """Get non-empty paragraph texts from a shape."""
    texts = []
    for para in shape.text_frame.paragraphs:
        t = para.text.strip()
        if t:
            texts.append(t)
    return texts


def has_centered_text(shape):
    """Check if shape has at least one centered paragraph with text."""
    for para in shape.text_frame.paragraphs:
        if para.text.strip():
            if para.alignment == PP_ALIGN.CENTER:
                return True
    return False


def has_bold_text(shape):
    """Check if shape has at least one bold run with text."""
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.text.strip() and run.font.bold:
                return True
    return False


def has_italic_text(shape):
    """Check if shape has at least one italic run with text."""
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.text.strip() and run.font.italic:
                return True
    return False


def get_font_color_rgb(run):
    """Safely get font color RGB or None."""
    try:
        if run.font.color.type is not None:
            return run.font.color.rgb
    except Exception:
        pass
    return None


def has_white_text(shape):
    """Check if shape has text with white color (FFFFFF)."""
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.text.strip():
                rgb = get_font_color_rgb(run)
                if rgb is not None and str(rgb) == 'FFFFFF':
                    return True
    return False


def get_slide_bg_color(slide):
    """Get slide background solid fill color or None."""
    fill = slide.background.fill
    if fill.type == 1:  # SOLID
        return fill.fore_color.rgb
    return None


def is_title_slide_pattern(slide):
    """Title Slide: centered large title (bold), centered subtitle, no solid bg."""
    text_shapes = get_text_shapes(slide)
    if len(text_shapes) < 2:
        return False
    # Need at least one centered bold text and one centered non-bold text
    has_title = False
    has_subtitle = False
    for shape in text_shapes:
        if has_centered_text(shape) and has_bold_text(shape):
            has_title = True
        elif has_centered_text(shape) and get_nonempty_text(shape):
            has_subtitle = True
    # No solid background for title slides
    bg = get_slide_bg_color(slide)
    if bg is not None:
        return False
    return has_title and has_subtitle


def is_content_layout_pattern(slide):
    """Content layout: has a rectangle/auto-shape (title bar) + text body area."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    has_rect = False
    has_body_text = False
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # Check if it's a title-bar-sized rectangle (wide, relatively short)
            if shape.width > Inches(5) and shape.height < Inches(2):
                has_rect = True
        if shape.has_text_frame:
            texts = get_nonempty_text(shape)
            # Body text: multiple lines of content text
            if len(texts) >= 3:
                has_body_text = True
    # No solid background expected
    bg = get_slide_bg_color(slide)
    if bg is not None:
        return False
    return has_rect and has_body_text


def is_two_column_pattern(slide):
    """Two Column: title + two side-by-side text boxes."""
    text_shapes = get_text_shapes(slide)
    if len(text_shapes) < 3:
        return False
    # Find shapes positioned side by side (similar top, different left)
    # Need two text shapes with similar vertical position but different horizontal
    shape_positions = []
    for shape in text_shapes:
        if get_nonempty_text(shape):
            shape_positions.append((shape.left, shape.top, shape.width, shape))

    # Look for pair of shapes with similar top and combined width fills the slide
    side_by_side = 0
    for i in range(len(shape_positions)):
        for j in range(i + 1, len(shape_positions)):
            l1, t1, w1, s1 = shape_positions[i]
            l2, t2, w2, s2 = shape_positions[j]
            # Similar top position (within 1 inch)
            if abs(t1 - t2) < Inches(1):
                # Different left positions (at least 2 inches apart)
                if abs(l1 - l2) > Inches(2):
                    side_by_side += 1
    return side_by_side >= 1


def is_image_focus_pattern(slide):
    """Image Focus: large rectangle placeholder + caption text."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    has_large_rect = False
    has_caption = False
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # Large rectangle: takes significant slide area
            if shape.width > Inches(5) and shape.height > Inches(3):
                has_large_rect = True
        if shape.has_text_frame:
            texts = get_nonempty_text(shape)
            if len(texts) >= 1 and has_italic_text(shape):
                has_caption = True
            elif len(texts) >= 1 and has_centered_text(shape):
                # Caption may not be italic but should be small centered text
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip() and run.font.size and run.font.size < Emu(254000):
                            has_caption = True
    return has_large_rect and has_caption


def is_quote_pattern(slide):
    """Quote: solid background, large italic centered text, attribution."""
    bg = get_slide_bg_color(slide)
    if bg is None:
        return False
    has_italic_centered = False
    has_attribution = False
    for shape in get_text_shapes(slide):
        texts = get_nonempty_text(shape)
        if texts:
            # Check for italic centered text (the quote)
            if has_italic_text(shape) and has_centered_text(shape):
                has_italic_centered = True
            # Check for attribution (text starting with dash/em-dash)
            for t in texts:
                if t.startswith('—') or t.startswith('-') or t.startswith('–'):
                    has_attribution = True
    return has_italic_centered and has_attribution


def is_section_divider_pattern(slide):
    """Section Divider: solid color background, centered white bold text."""
    bg = get_slide_bg_color(slide)
    if bg is None:
        return False
    text_shapes = get_text_shapes(slide)
    if len(text_shapes) < 1:
        return False
    # Need centered, white, bold text
    for shape in text_shapes:
        if has_centered_text(shape) and has_white_text(shape) and has_bold_text(shape):
            return True
    return False


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)
    num_slides = len(slides)

    # Component 1: Slide count >= 12 (0.20 points)
    # Initial has 1 slide; golden has 12. This differentiates.
    try:
        if num_slides >= 12:
            print(f"PASS: Component 1 — Slide count is {num_slides} (>= 12) (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 — Expected >= 12 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Title Slide layout (0.20 points)
    # At least 2 slides matching Title Slide pattern (centered bold title + subtitle)
    try:
        title_slides = [s for s in slides if is_title_slide_pattern(s)]
        count = len(title_slides)
        if count >= 2:
            print(f"PASS: Component 2 — Found {count} Title Slide layouts (>= 2) (0.20 pts)")
            total_score += 0.20
        elif count == 1:
            print(f"PARTIAL: Component 2 — Found {count} Title Slide layout (need 2) (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 — Found {count} Title Slide layouts (need >= 2)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Content layout (0.15 points)
    # At least 2 slides with title bar + body text area
    try:
        content_slides = [s for s in slides if is_content_layout_pattern(s)]
        count = len(content_slides)
        if count >= 2:
            print(f"PASS: Component 3 — Found {count} Content layouts (>= 2) (0.15 pts)")
            total_score += 0.15
        elif count == 1:
            print(f"PARTIAL: Component 3 — Found {count} Content layout (need 2) (0.075 pts)")
            total_score += 0.075
        else:
            print(f"FAIL: Component 3 — Found {count} Content layouts (need >= 2)")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Two Column layout (0.15 points)
    # At least 2 slides with title + two side-by-side text areas
    try:
        two_col_slides = [s for s in slides if is_two_column_pattern(s)]
        count = len(two_col_slides)
        if count >= 2:
            print(f"PASS: Component 4 — Found {count} Two Column layouts (>= 2) (0.15 pts)")
            total_score += 0.15
        elif count == 1:
            print(f"PARTIAL: Component 4 — Found {count} Two Column layout (need 2) (0.075 pts)")
            total_score += 0.075
        else:
            print(f"FAIL: Component 4 — Found {count} Two Column layouts (need >= 2)")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Image Focus layout (0.10 points)
    # At least 2 slides with large placeholder + caption
    try:
        image_slides = [s for s in slides if is_image_focus_pattern(s)]
        count = len(image_slides)
        if count >= 2:
            print(f"PASS: Component 5 — Found {count} Image Focus layouts (>= 2) (0.10 pts)")
            total_score += 0.10
        elif count == 1:
            print(f"PARTIAL: Component 5 — Found {count} Image Focus layout (need 2) (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 5 — Found {count} Image Focus layouts (need >= 2)")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Quote layout (0.10 points)
    # At least 2 slides with solid bg, italic centered text, and attribution
    try:
        quote_slides = [s for s in slides if is_quote_pattern(s)]
        count = len(quote_slides)
        if count >= 2:
            print(f"PASS: Component 6 — Found {count} Quote layouts (>= 2) (0.10 pts)")
            total_score += 0.10
        elif count == 1:
            print(f"PARTIAL: Component 6 — Found {count} Quote layout (need 2) (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 6 — Found {count} Quote layouts (need >= 2)")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Section Divider layout (0.10 points)
    # At least 2 slides with solid color bg + centered white bold text
    try:
        divider_slides = [s for s in slides if is_section_divider_pattern(s)]
        count = len(divider_slides)
        if count >= 2:
            print(f"PASS: Component 7 — Found {count} Section Divider layouts (>= 2) (0.10 pts)")
            total_score += 0.10
        elif count == 1:
            print(f"PARTIAL: Component 7 — Found {count} Section Divider layout (need 2) (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 7 — Found {count} Section Divider layouts (need >= 2)")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for unsaved GUI edits
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    # Also check Desktop for Template_Library.pptx as task instruction specifies
    alt_path = f'{WORKDIR}/Desktop/Template_Library.pptx'
    if os.path.exists(alt_path):
        file_path = alt_path
    else:
        print(f"File not found: {file_path} or {alt_path}")
        print("REWARD: 0.0")
        exit()

verify_task(file_path)
