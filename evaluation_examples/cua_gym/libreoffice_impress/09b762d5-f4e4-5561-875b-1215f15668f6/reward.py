"""
Reward Script: Process diagram on slide 7 with four labeled rectangles,
arrow connectors, alternating fill colors, curved return arrow, and white 14pt bold text.
Task ID: impress_ndo_065
Domain: libreoffice_impress
Scoring:
  Component 1: Four rectangles with labels Plan/Design/Build/Launch (0.30)
  Component 2: Alternating fill colors #3498DB/#2ECC71 (0.20)
  Component 3: Three arrow connectors between consecutive boxes (0.15)
  Component 4: Curved arrow from Launch back to Plan (0.15)
  Component 5: Text white 14pt bold (0.20)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_065'

NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}


def _check_connector_has_arrow(ln_element):
    """Check if a line element has an arrowhead on either end."""
    if ln_element is None:
        return False
    tail = ln_element.find('a:tailEnd', NS)
    head = ln_element.find('a:headEnd', NS)
    tail_has = tail is not None and tail.get('type', 'none') != 'none'
    head_has = head is not None and head.get('type', 'none') != 'none'
    return tail_has or head_has


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Pt, Emu
        from pptx.dml.color import RGBColor
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: at least 7 slides
    if len(prs.slides) < 7:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 7")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[6]  # Slide 7 (0-indexed)

    # Collect rectangles (AUTO_SHAPE with text) and connectors (LINE / FREEFORM)
    rectangles = []
    connectors_line = []

    for shape in slide.shapes:
        st = shape.shape_type
        # AUTO_SHAPE = 1
        if st == 1 and hasattr(shape, 'text') and shape.text.strip():
            rectangles.append(shape)
        # LINE = 9 (connectors show as LINE type in python-pptx)
        elif st == 9:
            connectors_line.append(shape)

    # ---------------------------------------------------------------
    # Component 1: Four rectangles with labels Plan/Design/Build/Launch (0.30 pts)
    # ---------------------------------------------------------------
    try:
        expected_labels = ['Plan', 'Design', 'Build', 'Launch']
        found_labels = [r.text.strip() for r in rectangles]
        # Check all 4 labels are present
        matched_labels = [lbl for lbl in expected_labels if lbl in found_labels]
        if len(matched_labels) == 4:
            print(f"PASS: Component 1 -- All 4 rectangles found with correct labels: {found_labels} (0.30 pts)")
            total_score += 0.30
        elif len(matched_labels) >= 2:
            partial = round(0.30 * len(matched_labels) / 4, 2)
            print(f"PARTIAL: Component 1 -- {len(matched_labels)}/4 labels found: {found_labels} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- Expected labels {expected_labels}, found rectangles with text: {found_labels}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Build a label->shape map for further checks
    label_shape_map = {}
    for r in rectangles:
        txt = r.text.strip()
        if txt in ['Plan', 'Design', 'Build', 'Launch']:
            label_shape_map[txt] = r

    # ---------------------------------------------------------------
    # Component 2: Alternating fill colors #3498DB and #2ECC71 (0.20 pts)
    # Plan=#3498DB, Design=#2ECC71, Build=#3498DB, Launch=#2ECC71
    # ---------------------------------------------------------------
    try:
        expected_colors = {
            'Plan': '3498DB',
            'Design': '2ECC71',
            'Build': '3498DB',
            'Launch': '2ECC71',
        }
        color_matches = 0
        total_to_check = 0
        for lbl, expected_hex in expected_colors.items():
            if lbl in label_shape_map:
                total_to_check += 1
                shape = label_shape_map[lbl]
                try:
                    fill = shape.fill
                    if fill.type is not None and fill.type == 1:  # SOLID
                        actual_hex = str(fill.fore_color.rgb).upper()
                        if actual_hex == expected_hex.upper():
                            color_matches += 1
                            print(f"  Color OK: {lbl} = #{actual_hex}")
                        else:
                            print(f"  Color MISMATCH: {lbl} expected #{expected_hex}, got #{actual_hex}")
                    else:
                        print(f"  Color FAIL: {lbl} fill is not solid (type={fill.type})")
                except Exception as e:
                    print(f"  Color ERROR: {lbl}: {e}")

        if total_to_check > 0 and color_matches == total_to_check:
            print(f"PASS: Component 2 -- All {color_matches} fill colors correct (0.20 pts)")
            total_score += 0.20
        elif color_matches > 0:
            partial = round(0.20 * color_matches / max(total_to_check, 4), 2)
            print(f"PARTIAL: Component 2 -- {color_matches}/{total_to_check} colors correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 -- No fill colors matched")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # ---------------------------------------------------------------
    # Component 3: Arrow connectors between consecutive boxes (0.15 pts)
    # Need at least 3 straight arrow connectors with arrowheads
    # ---------------------------------------------------------------
    try:
        # Parse XML to check for connector arrowheads and geometry
        straight_arrows = 0
        curved_arrows = 0

        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide7.xml') as f:
                tree = ET.parse(f)
                root = tree.getroot()
                for cxn in root.findall('.//p:cSld/p:spTree/p:cxnSp', NS):
                    spPr = cxn.find('p:spPr', NS)
                    if spPr is None:
                        continue
                    # Check geometry
                    prstGeom = spPr.find('a:prstGeom', NS)
                    geom = prstGeom.get('prst') if prstGeom is not None else None

                    # Check arrowhead on tail or head end
                    ln = spPr.find('a:ln', NS)
                    has_arrow = _check_connector_has_arrow(ln)

                    if geom == 'line' and has_arrow:
                        straight_arrows += 1
                    elif geom is not None and 'curved' in (geom or '').lower() and has_arrow:
                        curved_arrows += 1
                    elif geom is not None and geom != 'line' and has_arrow:
                        # Could be a non-straight connector (bentConnector, etc.)
                        curved_arrows += 1

        if straight_arrows >= 3:
            print(f"PASS: Component 3 -- {straight_arrows} straight arrow connectors found (0.15 pts)")
            total_score += 0.15
        elif straight_arrows >= 1:
            partial = round(0.15 * straight_arrows / 3, 2)
            print(f"PARTIAL: Component 3 -- {straight_arrows}/3 straight arrow connectors ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 -- No straight arrow connectors found")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # ---------------------------------------------------------------
    # Component 4: Curved arrow from Launch back to Plan (0.15 pts)
    # Must be a curved/non-straight connector with arrowhead
    # ---------------------------------------------------------------
    try:
        if curved_arrows >= 1:
            print(f"PASS: Component 4 -- Curved return arrow found ({curved_arrows} curved connector(s)) (0.15 pts)")
            total_score += 0.15
        else:
            # Also check for freeform shapes or arcs that might represent the curved arrow
            print(f"FAIL: Component 4 -- No curved arrow connector found (straight={straight_arrows}, curved={curved_arrows})")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # ---------------------------------------------------------------
    # Component 5: Text formatting - white (#FFFFFF), 14pt, bold (0.20 pts)
    # Check text in the 4 rectangle labels
    # ---------------------------------------------------------------
    try:
        format_checks = 0
        format_total = 0
        for lbl in ['Plan', 'Design', 'Build', 'Launch']:
            if lbl not in label_shape_map:
                continue
            shape = label_shape_map[lbl]
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not (run.text or '').strip():
                        continue
                    format_total += 1
                    checks_passed = 0
                    total_sub = 3  # bold, size, color

                    # Bold
                    if run.font.bold is True:
                        checks_passed += 1
                    else:
                        print(f"  Format FAIL: {lbl} bold={run.font.bold}")

                    # Size: 14pt = 177800 EMU
                    if run.font.size is not None and abs(run.font.size - 177800) < 1000:
                        checks_passed += 1
                    else:
                        print(f"  Format FAIL: {lbl} size={run.font.size} (expected 177800)")

                    # Color: white #FFFFFF
                    try:
                        if run.font.color.type is not None:
                            actual_rgb = str(run.font.color.rgb).upper()
                            if actual_rgb == 'FFFFFF':
                                checks_passed += 1
                            else:
                                print(f"  Format FAIL: {lbl} color=#{actual_rgb}")
                        else:
                            print(f"  Format FAIL: {lbl} color is inherited/None")
                    except Exception:
                        print(f"  Format FAIL: {lbl} color not accessible")

                    if checks_passed == total_sub:
                        format_checks += 1

        if format_total > 0 and format_checks == format_total:
            print(f"PASS: Component 5 -- All {format_checks} text runs have white 14pt bold (0.20 pts)")
            total_score += 0.20
        elif format_checks > 0:
            partial = round(0.20 * format_checks / max(format_total, 1), 2)
            print(f"PARTIAL: Component 5 -- {format_checks}/{format_total} text runs correctly formatted ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 -- No text runs with correct formatting")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
