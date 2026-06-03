"""
Initial Setup: Create Agile presentation with empty slide 7 for process diagram task.
Task ID: impress_ndo_065
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_065'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
    return slide


def add_title_only_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    add_title_slide(prs, "Agile Project Management", "A Comprehensive Overview for Modern Teams")

    # Slide 2: Agenda
    add_content_slide(prs, "Agenda", [
        "Introduction to Agile Methodology",
        "Core Principles and Values",
        "Sprint Planning and Execution",
        "Team Roles and Responsibilities",
        "Continuous Improvement Practices",
        "Development Cycle Overview",
        "Metrics and KPIs",
    ])

    # Slide 3: Core Principles
    add_content_slide(prs, "Core Agile Principles", [
        "Individuals and interactions over processes and tools",
        "Working software over comprehensive documentation",
        "Customer collaboration over contract negotiation",
        "Responding to change over following a plan",
    ])

    # Slide 4: Sprint Planning
    add_content_slide(prs, "Sprint Planning", [
        "Define sprint goals aligned with product backlog priorities",
        "Estimate story points using planning poker techniques",
        "Break epics into actionable user stories",
        "Assign tasks based on team capacity and velocity",
        "Set clear acceptance criteria for each story",
    ])

    # Slide 5: Team Roles
    add_content_slide(prs, "Team Roles", [
        "Product Owner: Manages backlog, defines priorities",
        "Scrum Master: Facilitates ceremonies, removes blockers",
        "Development Team: Cross-functional, self-organizing",
        "Stakeholders: Provide feedback during sprint reviews",
    ])

    # Slide 6: Continuous Improvement
    add_content_slide(prs, "Continuous Improvement", [
        "Conduct retrospectives after every sprint",
        "Track velocity trends over multiple sprints",
        "Identify and eliminate process bottlenecks",
        "Foster a culture of experimentation and learning",
        "Use data-driven decisions for process changes",
    ])

    # Slide 7: Development Cycle - EMPTY below title (task target)
    add_title_only_slide(prs, "Development Cycle")

    # Slide 8: Metrics
    add_content_slide(prs, "Key Metrics and KPIs", [
        "Sprint Velocity: Average story points per sprint",
        "Cycle Time: Time from start to completion",
        "Defect Rate: Bugs per release",
        "Customer Satisfaction: NPS score tracking",
        "Team Health: Happiness index surveys",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
