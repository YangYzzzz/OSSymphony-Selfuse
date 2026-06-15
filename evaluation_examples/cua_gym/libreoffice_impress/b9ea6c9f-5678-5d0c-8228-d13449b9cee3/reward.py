"""
Reward Script: Insert citation text box on slide 6
Task ID: impress_teach_013
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): Citation text box exists on slide 6 with correct text
  Component 2 (0.2): Font size is 10pt
  Component 3 (0.2): Font is italic
  Component 4 (0.2): Font color is #808080
  Component 5 (0.1): Text box is positioned near the bottom of the slide
"""

import os
from pptx import Presentation
from pptx.util import Pt, Emu

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_013'
EXPECTED_TEXT = 'Reference: Smith et al., 2023, Journal of Neuroscience, Vol. 45'


def persist_app_state(domain):
    """Save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            import time
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def find_citation_shape(slide):
    """Find the shape on slide 6 that contains the citation text.
    Returns the shape if found, None otherwise.
    """
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame') or not shape.has_text_frame:
            continue
        full_text = shape.text_frame.text.strip()
        if EXPECTED_TEXT.lower() in full_text.lower():
            return shape
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[5]  # 0-indexed, slide 6

    # Find the citation shape
    citation_shape = find_citation_shape(slide)

    # Component 1: Citation text box exists with correct text (0.3 points)
    try:
        if citation_shape is not None:
            actual_text = citation_shape.text_frame.text.strip()
            if actual_text == EXPECTED_TEXT:
                print(f"PASS: Component 1 -- Citation text matches exactly (0.3 pts)")
                total_score += 0.3
            else:
                # Partial: text contains the citation but not exact
                print(f"FAIL: Component 1 -- Text mismatch. Expected: {repr(EXPECTED_TEXT)}, Found: {repr(actual_text)}")
        else:
            print(f"FAIL: Component 1 -- No shape on slide 6 contains the citation text")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Components 2-5 only apply if citation shape was found
    if citation_shape is None:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Get runs from the citation shape for font checks
    runs = []
    for para in citation_shape.text_frame.paragraphs:
        for run in para.runs:
            if (run.text or "").strip():
                runs.append(run)

    if not runs:
        print("FAIL: Citation shape has no text runs to check formatting")
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Font size is 10pt / 127000 EMU (0.2 points)
    try:
        bad_size_runs = [r for r in runs if r.font.size is None or r.font.size != Pt(10)]
        if len(bad_size_runs) == 0:
            print(f"PASS: Component 2 -- Font size is 10pt ({Pt(10)} EMU) (0.2 pts)")
            total_score += 0.2
        else:
            r = bad_size_runs[0]
            print(f"FAIL: Component 2 -- Run '{r.text[:30]}...' has size {r.font.size} (expected {Pt(10)})")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Font is italic (0.2 points)
    try:
        # None means inherit (not italic), True means italic
        non_italic_runs = [r for r in runs if r.font.italic is not True]
        if len(non_italic_runs) == 0:
            print(f"PASS: Component 3 -- Font is italic (0.2 pts)")
            total_score += 0.2
        else:
            r = non_italic_runs[0]
            print(f"FAIL: Component 3 -- Run '{r.text[:30]}...' italic={r.font.italic} (expected True)")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Font color is #808080 (0.2 points)
    try:
        def get_run_color(run):
            try:
                if run.font.color.type is not None:
                    return str(run.font.color.rgb).upper()
            except Exception:
                pass
            return None

        bad_color_runs = [r for r in runs if get_run_color(r) != '808080']
        if len(bad_color_runs) == 0:
            print(f"PASS: Component 4 -- Font color is #808080 (0.2 pts)")
            total_score += 0.2
        else:
            r = bad_color_runs[0]
            print(f"FAIL: Component 4 -- Run '{r.text[:30]}...' color={get_run_color(r)} (expected 808080)")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Text box positioned near the bottom of the slide (0.1 points)
    # Slide height default is 6858000 EMU (7.5 inches). "Bottom" means top position
    # is in the lower portion of the slide. We check that top > 50% of slide height.
    try:
        slide_height = prs.slide_height
        shape_top = citation_shape.top
        # The text box should be in the bottom half of the slide
        if shape_top > slide_height * 0.5:
            print(f"PASS: Component 5 -- Text box at bottom (top={shape_top}, slide_height={slide_height}, ratio={shape_top/slide_height:.2f}) (0.1 pts)")
            total_score += 0.1
        else:
            print(f"FAIL: Component 5 -- Text box not at bottom (top={shape_top}, slide_height={slide_height}, ratio={shape_top/slide_height:.2f})")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = round(min(total_score, 1.0), 1)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
