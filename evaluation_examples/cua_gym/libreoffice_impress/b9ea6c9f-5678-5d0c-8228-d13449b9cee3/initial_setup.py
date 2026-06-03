"""
Initial Setup: Create a neuroscience lecture presentation with 8 slides.
Slide 6 has content about neural pathways but no citation.
Task ID: impress_teach_013
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_013'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_body(slide, title_text, body_lines, title_size=Pt(28), body_size=Pt(18)):
    """Helper to add title text box and body text box to a blank slide."""
    # Title text box
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.4), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = title_size
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    # Body text box
    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.0))
    btf = body_box.text_frame
    btf.word_wrap = True
    for i, line in enumerate(body_lines):
        if i == 0:
            p = btf.paragraphs[0]
        else:
            p = btf.add_paragraph()
        p.space_after = Pt(6)
        if line.strip():
            run = p.add_run()
            run.text = line
            run.font.size = body_size
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # Blank layout

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Fundamentals of Neuroscience"
    slide1.placeholders[1].text = "Dr. Emily Nakamura\nDepartment of Cognitive Sciences\nStanford University"

    # --- Slide 2: Course Overview ---
    slide2 = prs.slides.add_slide(blank_layout)
    add_title_and_body(slide2, "Course Overview", [
        "This lecture series explores the foundational principles of neuroscience,",
        "from cellular mechanisms to complex cognitive processes.",
        "",
        "Topics covered in this module:",
        "  - Neuronal structure and function",
        "  - Synaptic transmission and plasticity",
        "  - Neural circuit organization",
        "  - Sensory processing pathways",
        "  - Motor control and coordination",
    ])

    # --- Slide 3: Neuronal Structure ---
    slide3 = prs.slides.add_slide(blank_layout)
    add_title_and_body(slide3, "Neuronal Structure and Classification", [
        "Neurons are the fundamental computational units of the nervous system.",
        "The typical neuron consists of dendrites, a cell body (soma),",
        "and an axon that transmits signals to downstream targets.",
        "",
        "Classification by morphology:",
        "  - Multipolar neurons (most common in CNS)",
        "  - Bipolar neurons (retina, olfactory epithelium)",
        "  - Unipolar neurons (dorsal root ganglia)",
        "  - Pseudounipolar neurons (sensory afferents)",
    ])

    # --- Slide 4: Synaptic Transmission ---
    slide4 = prs.slides.add_slide(blank_layout)
    add_title_and_body(slide4, "Synaptic Transmission Mechanisms", [
        "Chemical synapses convert electrical signals to chemical messengers.",
        "Neurotransmitter release involves Ca2+-dependent vesicle fusion.",
        "",
        "Key neurotransmitter systems:",
        "  - Glutamate: primary excitatory transmitter",
        "  - GABA: primary inhibitory transmitter",
        "  - Dopamine: reward and motor control circuits",
        "  - Serotonin: mood regulation and sleep-wake cycles",
        "  - Acetylcholine: neuromuscular junction and attention",
    ])

    # --- Slide 5: Neural Plasticity ---
    slide5 = prs.slides.add_slide(blank_layout)
    add_title_and_body(slide5, "Neural Plasticity and Learning", [
        "Synaptic plasticity is the biological basis of learning and memory.",
        "Long-term potentiation (LTP) strengthens synaptic connections",
        "through repeated activation patterns.",
        "",
        "Forms of plasticity:",
        "  - Hebbian plasticity ('cells that fire together wire together')",
        "  - Spike-timing dependent plasticity (STDP)",
        "  - Homeostatic plasticity (synaptic scaling)",
        "  - Structural plasticity (dendritic spine remodeling)",
    ])

    # --- Slide 6: Neural Pathways (NO citation) ---
    slide6 = prs.slides.add_slide(blank_layout)
    add_title_and_body(slide6, "Neural Pathways and Signal Integration", [
        "Neural pathways form hierarchical processing streams in the brain.",
        "Ascending pathways carry sensory information to cortical areas,",
        "while descending pathways transmit motor commands to effectors.",
        "",
        "Major ascending pathways:",
        "  - Dorsal column-medial lemniscal pathway (fine touch, proprioception)",
        "  - Spinothalamic tract (pain, temperature)",
        "  - Visual pathway: retina -> LGN -> V1 -> ventral/dorsal streams",
        "",
        "Convergent and divergent connectivity patterns enable",
        "complex signal integration across cortical and subcortical regions.",
    ])

    # --- Slide 7: Motor Systems ---
    slide7 = prs.slides.add_slide(blank_layout)
    add_title_and_body(slide7, "Motor Control Systems", [
        "Voluntary movement requires coordinated activity across",
        "cortical, subcortical, and spinal motor circuits.",
        "",
        "Motor hierarchy:",
        "  - Primary motor cortex (M1): movement execution",
        "  - Premotor and supplementary motor areas: planning",
        "  - Basal ganglia: action selection and initiation",
        "  - Cerebellum: timing, coordination, error correction",
        "  - Spinal cord: pattern generators and reflexes",
    ])

    # --- Slide 8: Summary and Next Steps ---
    slide8 = prs.slides.add_slide(blank_layout)
    add_title_and_body(slide8, "Summary and Future Directions", [
        "This module covered the essential building blocks of neural computation:",
        "  - Neuronal structure and electrochemical signaling",
        "  - Synaptic transmission and neuromodulation",
        "  - Plasticity mechanisms underlying learning",
        "  - Sensory and motor pathway organization",
        "",
        "Next module: Advanced topics in computational neuroscience,",
        "including neural coding theories and brain-computer interfaces.",
        "",
        "Recommended reading: Kandel et al., Principles of Neural Science, 6th Ed.",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
