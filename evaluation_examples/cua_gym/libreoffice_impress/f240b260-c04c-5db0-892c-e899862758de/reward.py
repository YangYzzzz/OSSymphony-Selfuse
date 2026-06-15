"""
FINAL REWARD SCRIPT - SUCCESS
Task: I'm trying to reduce the file size of my presentation for quicker sharing. How can I compress all the images down to 96 DPI within LibreOffice Impress?
Generated: 2025-08-07 09:22:37
Status: success
Model: o4-mini
Total Steps: 3
"""

#!/usr/bin/env python3
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io

def verify_image_compression(file_path):
    print(f"Starting verification for image compression in: {file_path}")
    total_score = 0.0
    max_score = 1.0

    # Requirement 1: File existence (0.2)
    if os.path.exists(file_path):
        print("✓ File exists (0.2)")
        total_score += 0.2
    else:
        print("✗ File not found (0.0)")
        _report_and_exit(total_score, max_score)

    # Requirement 2: Presentation loads (0.1)
    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded with {len(prs.slides)} slides (0.1)")
        total_score += 0.1
    except Exception as e:
        print(f"✗ Error loading presentation: {e} (0.0)")
        _report_and_exit(total_score, max_score)

    # Requirement 3: Identify images (0.1)
    images = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images.append((slide_idx, shape))
    if images:
        print(f"✓ Found {len(images)} image(s) in presentation (0.1)")
        total_score += 0.1
    else:
        print("✗ No images found (0.0)")
        _report_and_exit(total_score, max_score)

    # Requirement 4: Check each image DPI <= 96 (total 0.6)
    per_image_score = 0.6 / len(images)
    for idx, (slide_idx, shape) in enumerate(images, start=1):
        img = shape.image
        blob = img.blob
        try:
            pil_img = Image.open(io.BytesIO(blob))
            dpi = pil_img.info.get('dpi')
            if dpi and isinstance(dpi, tuple) and dpi[0] <= 96 and dpi[1] <= 96:
                print(f"  ✓ Image {idx} on slide {slide_idx} DPI {dpi} <= (96,96) (+{per_image_score:.3f})")
                total_score += per_image_score
            else:
                print(f"  ✗ Image {idx} on slide {slide_idx} DPI {dpi} not <= (96,96) (+0.0)")
        except Exception as e:
            print(f"  ✗ Error reading image {idx} on slide {slide_idx}: {e} (+0.0)")

    final_score = min(total_score, max_score)
    print(f"Total score breakdown: {total_score}/{max_score}")
    print(f"Final score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score


def _report_and_exit(score, max_score):
    print(f"Total score: {score}/{max_score}")
    print(f"REWARD: {score}")
    exit(0)

if __name__ == '__main__':
    file_path = '/home/user/im_trying_to_reduce_the_file_size_of_my_presentation_for_quicker_sharing_how_can_i_compress_all_the_.pptx'
    verify_image_compression(file_path)
