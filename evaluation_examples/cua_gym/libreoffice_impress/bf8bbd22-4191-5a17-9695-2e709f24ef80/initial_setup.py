"""
Initial Setup: Create a contact slide presentation with 8 slides.
Task ID: impress_design_080
Domain: libreoffice_impress

Slide 8 has only a title 'Get In Touch' — no circles or contact info yet.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_design_080'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_name="Calibri",
                font_size=18, bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a styled textbox."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    add_textbox(slide1, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
                "Studio Design Co.", "Calibri", 44, True,
                RGBColor(0xFF, 0xFF, 0xFF), PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(2.5), Inches(3.8), Inches(8), Inches(1.0),
                "Creative Solutions for Modern Brands", "Calibri", 22, False,
                RGBColor(0xBD, 0xC3, 0xC7), PP_ALIGN.CENTER)

    # --- Slide 2: About Us ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide2, Inches(1.0), Inches(0.5), Inches(6), Inches(1.0),
                "About Us", "Calibri", 36, True,
                RGBColor(0x2C, 0x3E, 0x50))
    add_textbox(slide2, Inches(1.0), Inches(1.8), Inches(10), Inches(3.0),
                "Studio Design Co. is a full-service creative agency founded in 2018. "
                "We specialize in brand identity, digital marketing, and user experience design. "
                "Our team of 25 designers, strategists, and developers has delivered over 200 projects "
                "for clients ranging from startups to Fortune 500 companies.",
                "Calibri", 16, False, RGBColor(0x33, 0x33, 0x33))

    # --- Slide 3: Our Services ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide3, Inches(1.0), Inches(0.5), Inches(6), Inches(1.0),
                "Our Services", "Calibri", 36, True,
                RGBColor(0x2C, 0x3E, 0x50))
    services = [
        "Brand Identity & Logo Design",
        "Web Design & Development",
        "Mobile App UI/UX",
        "Social Media Strategy",
        "Print & Packaging Design",
    ]
    for i, svc in enumerate(services):
        add_textbox(slide3, Inches(1.5), Inches(1.8 + i * 0.7), Inches(8), Inches(0.6),
                    f"• {svc}", "Calibri", 18, False, RGBColor(0x33, 0x33, 0x33))

    # --- Slide 4: Portfolio Highlights ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide4, Inches(1.0), Inches(0.5), Inches(6), Inches(1.0),
                "Portfolio Highlights", "Calibri", 36, True,
                RGBColor(0x2C, 0x3E, 0x50))
    projects = [
        ("Meridian Hotels", "Complete rebrand including logo, website, and collateral"),
        ("TechFlow App", "Mobile banking app serving 500K+ users"),
        ("Bloom Cosmetics", "E-commerce platform with 40% conversion increase"),
    ]
    for i, (name, desc) in enumerate(projects):
        add_textbox(slide4, Inches(1.5), Inches(1.8 + i * 1.2), Inches(9), Inches(0.5),
                    name, "Calibri", 22, True, RGBColor(0x2C, 0x3E, 0x50))
        add_textbox(slide4, Inches(1.5), Inches(2.3 + i * 1.2), Inches(9), Inches(0.5),
                    desc, "Calibri", 14, False, RGBColor(0x7F, 0x8C, 0x8D))

    # --- Slide 5: Our Process ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide5, Inches(1.0), Inches(0.5), Inches(6), Inches(1.0),
                "Our Process", "Calibri", 36, True,
                RGBColor(0x2C, 0x3E, 0x50))
    steps = ["Discovery & Research", "Strategy & Planning",
             "Design & Prototyping", "Development & Testing", "Launch & Support"]
    for i, step in enumerate(steps):
        add_textbox(slide5, Inches(1.5), Inches(1.8 + i * 0.7), Inches(8), Inches(0.6),
                    f"{i + 1}. {step}", "Calibri", 18, False, RGBColor(0x33, 0x33, 0x33))

    # --- Slide 6: Team ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide6, Inches(1.0), Inches(0.5), Inches(6), Inches(1.0),
                "Meet the Team", "Calibri", 36, True,
                RGBColor(0x2C, 0x3E, 0x50))
    team = [
        ("Elena Rodriguez", "Creative Director", "15 years in brand design"),
        ("James Park", "Lead Developer", "Full-stack with React & Node.js"),
        ("Aisha Patel", "UX Strategist", "Former Google UX researcher"),
    ]
    for i, (name, role, bio) in enumerate(team):
        add_textbox(slide6, Inches(1.5), Inches(1.8 + i * 1.3), Inches(4), Inches(0.5),
                    name, "Calibri", 20, True, RGBColor(0x2C, 0x3E, 0x50))
        add_textbox(slide6, Inches(1.5), Inches(2.2 + i * 1.3), Inches(4), Inches(0.4),
                    role, "Calibri", 16, False, RGBColor(0x34, 0x98, 0xDB))
        add_textbox(slide6, Inches(1.5), Inches(2.55 + i * 1.3), Inches(6), Inches(0.4),
                    bio, "Calibri", 14, False, RGBColor(0x7F, 0x8C, 0x8D))

    # --- Slide 7: Testimonials ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide7, Inches(1.0), Inches(0.5), Inches(6), Inches(1.0),
                "Client Testimonials", "Calibri", 36, True,
                RGBColor(0x2C, 0x3E, 0x50))
    add_textbox(slide7, Inches(1.5), Inches(2.0), Inches(9), Inches(1.5),
                '"Studio Design transformed our brand presence. Their attention to detail '
                'and creative vision exceeded our expectations." — Sarah Kim, CEO of Bloom Cosmetics',
                "Calibri", 16, False, RGBColor(0x55, 0x55, 0x55))
    add_textbox(slide7, Inches(1.5), Inches(4.0), Inches(9), Inches(1.5),
                '"The team delivered a world-class mobile app on time and under budget. '
                'Highly recommend for any digital project." — David Chen, CTO of TechFlow',
                "Calibri", 16, False, RGBColor(0x55, 0x55, 0x55))

    # --- Slide 8: Get In Touch (INITIAL — title only, NO contact info) ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide8.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    add_textbox(slide8, Inches(1.0), Inches(0.8), Inches(11), Inches(1.2),
                "Get In Touch", "Calibri", 40, True,
                RGBColor(0x2C, 0x3E, 0x50), PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
