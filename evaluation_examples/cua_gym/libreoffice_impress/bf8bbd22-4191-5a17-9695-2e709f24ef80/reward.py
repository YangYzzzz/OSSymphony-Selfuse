"""
Reward Script: Contact information slide with colored circles and text
Task ID: impress_design_080
Domain: libreoffice_impress
Scoring:
  - 3 colored circles (email/phone/web) at correct positions with correct colors (0.15 each = 0.45)
  - 3 text boxes with correct content, font, size, color (0.20 + 0.20 + 0.15 = 0.55)
  Total: 1.0
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_design_080'

# Expected data for the 3 contact rows
CONTACT_ROWS = [
    {
        "label": "Email",
        "circle_color": "3498DB",
        "circle_y_in": 2.5,
        "text": "hello@studio.design",
        "circle_pts": 0.15,
        "text_pts": 0.20,
    },
    {
        "label": "Phone",
        "circle_color": "2ECC71",
        "circle_y_in": 3.5,
        "text": "+1 (555) 123-4567",
        "circle_pts": 0.15,
        "text_pts": 0.20,
    },
    {
        "label": "Web",
        "circle_color": "E74C3C",
        "circle_y_in": 4.5,
        "text": "www.studio.design",
        "circle_pts": 0.15,
        "text_pts": 0.15,
    },
]

CIRCLE_X_IN = 2.0
CIRCLE_DIAMETER_IN = 0.3
TEXT_X_IN = 2.8
TEXT_FONT = "Roboto"
TEXT_SIZE = Pt(18)  # 228600 EMU
TEXT_COLOR = "333333"

# Tolerance for position/size checks (in EMU)
POS_TOLERANCE = Inches(0.15)  # generous tolerance for position
SIZE_TOLERANCE = Inches(0.1)


def approx_eq(val, expected, tolerance):
    """Check if val is approximately equal to expected within tolerance (EMU)."""
    return abs(val - expected) <= tolerance


def find_circles_on_slide(slide):
    """Find all oval/auto shapes that look like small circles on the slide."""
    circles = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # Check if it's roughly circular (width ~= height, small)
            if shape.width > 0 and shape.height > 0:
                try:
                    fill = shape.fill
                    if fill.type is not None:
                        color = str(fill.fore_color.rgb)
                    else:
                        color = None
                except:
                    color = None
                circles.append({
                    "shape": shape,
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                    "color": color,
                })
    return circles


def find_text_boxes_on_slide(slide):
    """Find all text box shapes (excluding the title)."""
    boxes = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shape.has_text_frame:
            text = ""
            font_name = None
            font_size = None
            font_color = None
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    text += run.text
                    if font_name is None:
                        font_name = run.font.name
                    if font_size is None:
                        font_size = run.font.size
                    if font_color is None:
                        try:
                            if run.font.color.type is not None:
                                font_color = str(run.font.color.rgb)
                        except:
                            pass
            boxes.append({
                "shape": shape,
                "left": shape.left,
                "top": shape.top,
                "text": text.strip(),
                "font_name": font_name,
                "font_size": font_size,
                "font_color": font_color,
            })
    return boxes


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 8 slides
    if len(prs.slides) < 8:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 8")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[7]  # Slide 8 (0-indexed)

    # Find all circles and text boxes on slide 8
    circles = find_circles_on_slide(slide)
    text_boxes = find_text_boxes_on_slide(slide)

    # Filter text boxes: exclude the title "Get In Touch"
    contact_text_boxes = [tb for tb in text_boxes if tb["text"].lower() != "get in touch"]

    print(f"Found {len(circles)} circle shapes and {len(contact_text_boxes)} non-title text boxes on slide 8")

    for row in CONTACT_ROWS:
        label = row["label"]
        expected_color = row["circle_color"]
        expected_y = Inches(row["circle_y_in"])
        expected_text = row["text"]
        circle_pts = row["circle_pts"]
        text_pts = row["text_pts"]

        # --- Component: Circle check ---
        try:
            # Find a circle with matching color near the expected position
            matched_circle = None
            for c in circles:
                color_match = (c["color"] is not None and
                               c["color"].upper() == expected_color.upper())
                y_match = approx_eq(c["top"], expected_y, POS_TOLERANCE)
                x_match = approx_eq(c["left"], Inches(CIRCLE_X_IN), POS_TOLERANCE)
                size_match = (approx_eq(c["width"], Inches(CIRCLE_DIAMETER_IN), SIZE_TOLERANCE) and
                              approx_eq(c["height"], Inches(CIRCLE_DIAMETER_IN), SIZE_TOLERANCE))

                if color_match and y_match:
                    matched_circle = c
                    break

            if matched_circle is not None:
                # Verify position and size details
                details = []
                if not approx_eq(matched_circle["left"], Inches(CIRCLE_X_IN), POS_TOLERANCE):
                    details.append(f"x={matched_circle['left']/914400:.2f}in (expected {CIRCLE_X_IN}in)")
                if not (approx_eq(matched_circle["width"], Inches(CIRCLE_DIAMETER_IN), SIZE_TOLERANCE) and
                        approx_eq(matched_circle["height"], Inches(CIRCLE_DIAMETER_IN), SIZE_TOLERANCE)):
                    details.append(f"size={matched_circle['width']/914400:.2f}x{matched_circle['height']/914400:.2f}in")

                if details:
                    print(f"PASS (partial): {label} circle found with color #{expected_color} at y={row['circle_y_in']}in, but {'; '.join(details)} ({circle_pts} pts)")
                else:
                    print(f"PASS: {label} circle - color #{expected_color}, pos correct, size 0.3in ({circle_pts} pts)")
                if matched_circle is not None:
                    total_score += circle_pts
            else:
                print(f"FAIL: {label} circle - no oval with color #{expected_color} found near y={row['circle_y_in']}in")
        except Exception as e:
            print(f"ERROR: {label} circle check - {e}")

        # --- Component: Text check ---
        try:
            matched_text = None
            for tb in contact_text_boxes:
                if expected_text.lower() in tb["text"].lower():
                    matched_text = tb
                    break

            if matched_text is not None:
                sub_score = 0.0
                max_sub = text_pts

                # Check text content (40% of text points)
                content_pts = max_sub * 0.4
                if matched_text["text"].strip() == expected_text:
                    print(f"  PASS: {label} text content = '{expected_text}'")
                    sub_score += content_pts
                elif expected_text.lower() in matched_text["text"].lower():
                    print(f"  PARTIAL: {label} text contains '{expected_text}' in '{matched_text['text']}'")
                    sub_score += content_pts * 0.5

                # Check font name (20% of text points)
                font_pts = max_sub * 0.2
                if matched_text["font_name"] is not None and matched_text["font_name"] == TEXT_FONT:
                    print(f"  PASS: {label} font = {TEXT_FONT}")
                    sub_score += font_pts
                else:
                    print(f"  FAIL: {label} font = {matched_text['font_name']} (expected {TEXT_FONT})")

                # Check font size (20% of text points)
                size_pts = max_sub * 0.2
                if matched_text["font_size"] is not None and abs(matched_text["font_size"] - TEXT_SIZE) < Pt(2):
                    print(f"  PASS: {label} size = 18pt")
                    sub_score += size_pts
                else:
                    actual_pt = matched_text["font_size"] / 12700 if matched_text["font_size"] else None
                    print(f"  FAIL: {label} size = {actual_pt}pt (expected 18pt)")

                # Check font color (20% of text points)
                color_pts = max_sub * 0.2
                if matched_text["font_color"] is not None and matched_text["font_color"].upper() == TEXT_COLOR.upper():
                    print(f"  PASS: {label} color = #{TEXT_COLOR}")
                    sub_score += color_pts
                else:
                    print(f"  FAIL: {label} color = {matched_text['font_color']} (expected #{TEXT_COLOR})")

                # Check x position
                if approx_eq(matched_text["left"], Inches(TEXT_X_IN), POS_TOLERANCE):
                    print(f"  PASS: {label} text x position ~= {TEXT_X_IN}in")
                else:
                    print(f"  INFO: {label} text x = {matched_text['left']/914400:.2f}in (expected {TEXT_X_IN}in)")

                if sub_score > 0:
                    print(f"PASS: {label} text - scored {sub_score:.3f}/{max_sub} pts")
                    total_score += sub_score
                else:
                    print(f"FAIL: {label} text - all sub-checks failed")
            else:
                print(f"FAIL: {label} text - '{expected_text}' not found on slide 8")
        except Exception as e:
            print(f"ERROR: {label} text check - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.3f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
