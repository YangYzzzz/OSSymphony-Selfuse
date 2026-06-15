"""
Reward Script: Apply underline and dark red (#8B0000) to all content text on slide 5
Task ID: osworld_impress_underline_darkred_table_005
Domain: libreoffice_impress
Scoring:
  Component 1: Bullet textbox (4 items) — underline=True AND color=8B0000  (0.4 pts)
  Component 2: Table (4x2, all 8 cells) — underline=True AND color=8B0000  (0.4 pts)
  Component 3: Caption textbox         — underline=True AND color=8B0000  (0.2 pts)
Total: 1.0
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_underline_darkred_table_005'
TARGET_COLOR = '8B0000'
SLIDE_IDX = 4  # Slide 5 is index 4 (0-based)


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Validate slide count
    if len(prs.slides) < 5:
        print(f"CRITICAL: Expected at least 5 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[SLIDE_IDX]

    # Identify the relevant shapes on slide 5
    bullet_textbox = None   # TextBox 3 — 4 bullet items in black initially
    table_shape = None      # Table 4 — 4x2 table
    caption_textbox = None  # TextBox 5 — caption below table

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_shape = shape
        elif shape.has_text_frame:
            # Skip the header textbox (name="TextBox 2") which is NOT a content shape
            if shape.name == 'TextBox 2':
                continue
            # Check how many paragraphs with content — bullet textbox has 4 items
            paras_with_text = [p for p in shape.text_frame.paragraphs
                               if any((r.text or '').strip() for r in p.runs)]
            if len(paras_with_text) >= 4:
                bullet_textbox = shape
            elif shape.name.startswith('TextBox'):
                # Caption has 1 paragraph
                caption_textbox = shape

    # -------------------------------------------------------------------
    # Component 1: Bullet textbox — 4 items underline=True and color=8B0000 (0.4 pts)
    # -------------------------------------------------------------------
    try:
        if bullet_textbox is None:
            print("FAIL: Component 1 — Bullet textbox not found on slide 5")
        else:
            paras_ok = 0
            total_paras = 0
            for para in bullet_textbox.text_frame.paragraphs:
                runs_with_text = [r for r in para.runs if (r.text or '').strip()]
                if not runs_with_text:
                    continue
                total_paras += 1
                para_ok = True
                for run in runs_with_text:
                    # Check underline
                    if run.font.underline is not True:
                        para_ok = False
                        print(f"FAIL: Component 1 — bullet para '{para.text[:30]}' underline={run.font.underline}, expected True")
                        break
                    # Check color
                    try:
                        if run.font.color.type is None:
                            para_ok = False
                            print(f"FAIL: Component 1 — bullet para '{para.text[:30]}' has no explicit color")
                            break
                        actual_color = str(run.font.color.rgb).upper()
                        if actual_color != TARGET_COLOR.upper():
                            para_ok = False
                            print(f"FAIL: Component 1 — bullet para '{para.text[:30]}' color={actual_color}, expected {TARGET_COLOR}")
                            break
                    except Exception as ce:
                        para_ok = False
                        print(f"FAIL: Component 1 — color check error: {ce}")
                        break
                if para_ok:
                    paras_ok += 1

            if total_paras >= 4 and paras_ok == total_paras:
                print(f"PASS: Component 1 — All {paras_ok} bullet items have underline=True and color=8B0000 (0.4 pts)")
                total_score += 0.4
            else:
                print(f"FAIL: Component 1 — {paras_ok}/{total_paras} bullet items correctly formatted")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # -------------------------------------------------------------------
    # Component 2: Table cells — all 8 cells underline=True and color=8B0000 (0.4 pts)
    # -------------------------------------------------------------------
    try:
        if table_shape is None:
            print("FAIL: Component 2 — Table not found on slide 5")
        else:
            table = table_shape.table
            num_rows = len(table.rows)
            num_cols = len(table.columns)
            cells_ok = 0
            total_cells = 0
            for r in range(num_rows):
                for c in range(num_cols):
                    cell = table.cell(r, c)
                    runs_with_text = [run
                                      for para in cell.text_frame.paragraphs
                                      for run in para.runs
                                      if (run.text or '').strip()]
                    if not runs_with_text:
                        continue
                    total_cells += 1
                    cell_ok = True
                    for run in runs_with_text:
                        # Check underline
                        if run.font.underline is not True:
                            cell_ok = False
                            print(f"FAIL: Component 2 — cell[{r},{c}] underline={run.font.underline}, expected True")
                            break
                        # Check color
                        try:
                            if run.font.color.type is None:
                                cell_ok = False
                                print(f"FAIL: Component 2 — cell[{r},{c}] has no explicit color")
                                break
                            actual_color = str(run.font.color.rgb).upper()
                            if actual_color != TARGET_COLOR.upper():
                                cell_ok = False
                                print(f"FAIL: Component 2 — cell[{r},{c}] color={actual_color}, expected {TARGET_COLOR}")
                                break
                        except Exception as ce:
                            cell_ok = False
                            print(f"FAIL: Component 2 — cell[{r},{c}] color check error: {ce}")
                            break
                    if cell_ok:
                        cells_ok += 1

            if total_cells >= 8 and cells_ok == total_cells:
                print(f"PASS: Component 2 — All {cells_ok} table cells have underline=True and color=8B0000 (0.4 pts)")
                total_score += 0.4
            else:
                print(f"FAIL: Component 2 — {cells_ok}/{total_cells} table cells correctly formatted")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # -------------------------------------------------------------------
    # Component 3: Caption textbox — underline=True and color=8B0000 (0.2 pts)
    # -------------------------------------------------------------------
    try:
        if caption_textbox is None:
            print("FAIL: Component 3 — Caption textbox not found on slide 5")
        else:
            runs_with_text = [run
                              for para in caption_textbox.text_frame.paragraphs
                              for run in para.runs
                              if (run.text or '').strip()]
            if not runs_with_text:
                print("FAIL: Component 3 — Caption textbox has no text runs")
            else:
                caption_ok = True
                for run in runs_with_text:
                    if run.font.underline is not True:
                        caption_ok = False
                        print(f"FAIL: Component 3 — caption underline={run.font.underline}, expected True")
                        break
                    try:
                        if run.font.color.type is None:
                            caption_ok = False
                            print("FAIL: Component 3 — caption has no explicit color")
                            break
                        actual_color = str(run.font.color.rgb).upper()
                        if actual_color != TARGET_COLOR.upper():
                            caption_ok = False
                            print(f"FAIL: Component 3 — caption color={actual_color}, expected {TARGET_COLOR}")
                            break
                    except Exception as ce:
                        caption_ok = False
                        print(f"FAIL: Component 3 — caption color check error: {ce}")
                        break
                if caption_ok:
                    print(f"PASS: Component 3 — Caption textbox has underline=True and color=8B0000 (0.2 pts)")
                    total_score += 0.2
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
