"""
Initial Setup: 7-slide quarterly review deck — slide 5 has body textbox (4 items),
               4x2 metrics table, and caption textbox, all in plain black text.
Task ID: osworld_impress_underline_darkred_table_005
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_underline_darkred_table_005'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, text, left, top, width, height, font_size=16, bold=False, alignment=PP_ALIGN.LEFT):
    """Helper: add a plain black textbox."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.underline = False
    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black
    return txBox


def add_title(slide, title_text, subtitle_text=None):
    """Set title placeholder text."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text
        slide.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    if subtitle_text and len(slide.placeholders) > 1:
        try:
            ph = slide.placeholders[1]
            ph.text = subtitle_text
            if ph.text_frame.paragraphs:
                for run in ph.text_frame.paragraphs[0].runs:
                    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        except Exception:
            pass


def create_initial():
    prs = Presentation()
    # Use standard slide size (widescreen 13.33 x 7.5 in)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_layouts = prs.slide_layouts
    # Layout 0: Title Slide, Layout 1: Title+Content, Layout 5: Blank, Layout 6: Title Only

    # ------------------------------------------------------------------ #
    # Slide 1: Title slide — "Q3 2025 Quarterly Business Review"          #
    # ------------------------------------------------------------------ #
    slide1 = prs.slides.add_slide(slide_layouts[0])
    tf_title = slide1.shapes.title.text_frame
    tf_title.paragraphs[0].text = "Q3 2025 Quarterly Business Review"
    for run in tf_title.paragraphs[0].runs:
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
    try:
        ph_sub = slide1.placeholders[1]
        ph_sub.text = "Presented by Finance & Strategy Team\nSeptember 30, 2025"
        for para in ph_sub.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0x40, 0x40, 0x40)
    except Exception:
        pass

    # ------------------------------------------------------------------ #
    # Slide 2: Executive Summary                                           #
    # ------------------------------------------------------------------ #
    slide2 = prs.slides.add_slide(slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    for run in slide2.shapes.title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
    try:
        content_ph = slide2.placeholders[1]
        content_ph.text = (
            "• Total revenue reached $4.8M, up 12% YoY\n"
            "• Operating margin improved to 22.4%\n"
            "• New customer acquisition: 340 accounts\n"
            "• Product launches: 3 new SKUs shipped\n"
            "• Headcount grew from 128 to 145 employees"
        )
        for para in content_ph.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(0x20, 0x20, 0x20)
    except Exception:
        pass

    # ------------------------------------------------------------------ #
    # Slide 3: Revenue Breakdown                                           #
    # ------------------------------------------------------------------ #
    slide3 = prs.slides.add_slide(slide_layouts[5])  # blank
    # Title textbox
    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf3 = title3.text_frame
    p3 = tf3.paragraphs[0]
    r3 = p3.add_run()
    r3.text = "Revenue Breakdown by Region"
    r3.font.size = Pt(28)
    r3.font.bold = True
    r3.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Revenue table (5 rows x 3 cols)
    tbl3 = slide3.shapes.add_table(5, 3, Inches(1.5), Inches(1.3), Inches(10), Inches(3.5))
    t3 = tbl3.table
    headers3 = ["Region", "Q3 Revenue ($K)", "% of Total"]
    for c, h in enumerate(headers3):
        cell = t3.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
    rows3 = [
        ["North America", "2,340", "48.8%"],
        ["Europe",        "1,105", "23.0%"],
        ["Asia-Pacific",   "920", "19.2%"],
        ["Rest of World",  "435",  "9.0%"],
    ]
    for r, row_data in enumerate(rows3, 1):
        for c, val in enumerate(row_data):
            cell = t3.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(13)
                run.font.color.rgb = RGBColor(0x20, 0x20, 0x20)

    # ------------------------------------------------------------------ #
    # Slide 4: Operational Highlights                                      #
    # ------------------------------------------------------------------ #
    slide4 = prs.slides.add_slide(slide_layouts[5])  # blank
    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf4 = title4.text_frame
    p4 = tf4.paragraphs[0]
    r4 = p4.add_run()
    r4.text = "Operational Highlights"
    r4.font.size = Pt(28)
    r4.font.bold = True
    r4.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    bullets4 = [
        "Launched new CRM platform reducing support tickets by 18%",
        "Completed ISO 27001 recertification audit",
        "Expanded data center capacity by 35% in APAC region",
        "Employee NPS score rose from 42 to 61 this quarter",
        "Reduced average order fulfillment time from 3.2 to 2.1 days",
    ]
    body4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5))
    tf_b4 = body4.text_frame
    tf_b4.word_wrap = True
    for i, bullet_text in enumerate(bullets4):
        if i == 0:
            p = tf_b4.paragraphs[0]
        else:
            p = tf_b4.add_paragraph()
        run = p.add_run()
        run.text = f"• {bullet_text}"
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x20, 0x20, 0x20)

    # ------------------------------------------------------------------ #
    # Slide 5: Key Performance Metrics (THE TARGET SLIDE)                  #
    # ------------------------------------------------------------------ #
    slide5 = prs.slides.add_slide(slide_layouts[5])  # blank

    # Title
    title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.75))
    tf5_title = title5.text_frame
    p5_title = tf5_title.paragraphs[0]
    r5_title = p5_title.add_run()
    r5_title.text = "Key Performance Metrics — Q3 2025"
    r5_title.font.size = Pt(26)
    r5_title.font.bold = True
    r5_title.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Body textbox: 4 bullet items in plain BLACK text (NO underline, NO dark red)
    body5 = slide5.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(5.8), Inches(3.2))
    tf_body5 = body5.text_frame
    tf_body5.word_wrap = True
    bullet_items = [
        "Revenue growth exceeded forecast by 4.2%",
        "Customer churn rate held below 2.5%",
        "Net Promoter Score improved to 68",
        "R&D spend efficiency ratio reached 0.41",
    ]
    for i, item in enumerate(bullet_items):
        if i == 0:
            p = tf_body5.paragraphs[0]
        else:
            p = tf_body5.add_paragraph()
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(15)
        run.font.underline = False
        run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black

    # 4×2 metrics table (NO underline, NO dark red)
    tbl5 = slide5.shapes.add_table(4, 2, Inches(6.7), Inches(1.1), Inches(6.0), Inches(3.2))
    t5 = tbl5.table
    table_data = [
        ["Metric",                  "Value"],
        ["Total Revenue",           "$4.80M"],
        ["Operating Margin",        "22.4%"],
        ["New Customer Accounts",   "340"],
    ]
    for r_idx, row_data in enumerate(table_data):
        for c_idx, val in enumerate(row_data):
            cell = t5.cell(r_idx, c_idx)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(14)
                run.font.bold = (r_idx == 0)  # header row bold
                run.font.underline = False
                run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black

    # Caption textbox below the table (NO underline, NO dark red)
    caption5 = slide5.shapes.add_textbox(Inches(6.7), Inches(4.55), Inches(6.0), Inches(0.5))
    tf_caption5 = caption5.text_frame
    p_caption5 = tf_caption5.paragraphs[0]
    r_caption5 = p_caption5.add_run()
    r_caption5.text = "Source: Internal Finance Report — September 2025"
    r_caption5.font.size = Pt(11)
    r_caption5.font.italic = True
    r_caption5.font.underline = False
    r_caption5.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black

    # ------------------------------------------------------------------ #
    # Slide 6: Headcount & HR Summary                                      #
    # ------------------------------------------------------------------ #
    slide6 = prs.slides.add_slide(slide_layouts[5])  # blank
    title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    p6_t = title6.text_frame.paragraphs[0]
    r6_t = p6_t.add_run()
    r6_t.text = "Headcount & HR Summary"
    r6_t.font.size = Pt(28)
    r6_t.font.bold = True
    r6_t.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    hr_items = [
        "Total headcount: 145 FTEs (up from 128 in Q2)",
        "Engineering: 58 | Sales: 32 | Support: 25 | G&A: 30",
        "Open requisitions: 12 roles across 4 departments",
        "Average tenure: 2.8 years; voluntary attrition: 6.3%",
    ]
    body6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.5), Inches(4.5))
    tf_b6 = body6.text_frame
    tf_b6.word_wrap = True
    for i, item in enumerate(hr_items):
        if i == 0:
            p = tf_b6.paragraphs[0]
        else:
            p = tf_b6.add_paragraph()
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x20, 0x20, 0x20)

    # ------------------------------------------------------------------ #
    # Slide 7: Outlook & Next Steps                                        #
    # ------------------------------------------------------------------ #
    slide7 = prs.slides.add_slide(slide_layouts[5])  # blank
    title7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    p7_t = title7.text_frame.paragraphs[0]
    r7_t = p7_t.add_run()
    r7_t.text = "Q4 2025 Outlook & Next Steps"
    r7_t.font.size = Pt(28)
    r7_t.font.bold = True
    r7_t.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    outlook_items = [
        "Target revenue of $5.2M in Q4 (8% QoQ growth)",
        "Launch two additional SKUs by November 15",
        "Complete migration to cloud-native ERP by year-end",
        "Finalize 2026 strategic plan and board presentation",
        "Hire remaining 12 open roles before December 31",
    ]
    body7 = slide7.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.0))
    tf_b7 = body7.text_frame
    tf_b7.word_wrap = True
    for i, item in enumerate(outlook_items):
        if i == 0:
            p = tf_b7.paragraphs[0]
        else:
            p = tf_b7.add_paragraph()
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x20, 0x20, 0x20)

    # ------------------------------------------------------------------ #
    # Save                                                                  #
    # ------------------------------------------------------------------ #
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
