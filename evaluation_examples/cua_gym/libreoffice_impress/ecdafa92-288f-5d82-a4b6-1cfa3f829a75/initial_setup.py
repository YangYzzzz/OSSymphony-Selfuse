"""
Initial Setup: Export slide 1 as a JPEG image to the Desktop.
Task ID: impstruct_040
Domain: libreoffice_impress

Creates a 3-slide presentation 'cover_slide.pptx' and opens it in LibreOffice Impress.
No JPEG export exists yet — that is the agent's task.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impstruct_040'
OUTPUT = f'{WORKDIR}/cover_slide.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide (this is the one to be exported) ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q3 2025 Strategic Review"
    slide1.placeholders[1].text = "Prepared by the Corporate Strategy Division\nJuly 2025"

    # Set a distinct background color for slide 1
    fill1 = slide1.background.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)  # dark navy

    # Style the title text
    for run in slide1.shapes.title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.name = "Arial"

    # Style the subtitle
    for para in slide1.placeholders[1].text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
            run.font.size = Pt(18)
            run.font.name = "Arial"

    # --- Slide 2: Content Slide ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Revenue Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Total Revenue: $14.2M (up 12% YoY)"
    p2 = body2.add_paragraph()
    p2.text = "North America: $8.1M"
    p2.level = 1
    p3 = body2.add_paragraph()
    p3.text = "Europe: $3.9M"
    p3.level = 1
    p4 = body2.add_paragraph()
    p4.text = "Asia-Pacific: $2.2M"
    p4.level = 1
    p5 = body2.add_paragraph()
    p5.text = ""
    p6 = body2.add_paragraph()
    p6.text = "Key growth driver: Enterprise SaaS subscriptions (+23%)"

    # --- Slide 3: Summary Slide ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Next Steps & Action Items"
    body3 = slide3.placeholders[1].text_frame
    items = [
        "Finalize partnership agreement with Meridian Corp by Aug 15",
        "Launch Phase 2 marketing campaign targeting EMEA region",
        "Complete hiring for 5 senior engineering positions",
        "Present updated roadmap at September all-hands meeting",
        "Schedule quarterly review with board of directors",
    ]
    body3.text = items[0]
    for item in items[1:]:
        p = body3.add_paragraph()
        p.text = item

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Ensure Desktop directory exists (VM user is 'user', not root)
    os.makedirs('/home/user/Desktop', exist_ok=True)

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
