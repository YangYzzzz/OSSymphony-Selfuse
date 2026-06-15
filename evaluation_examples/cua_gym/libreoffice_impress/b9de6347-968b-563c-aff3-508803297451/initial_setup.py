"""
Initial Setup: Create Meeting.pptx with 6 slides, each with speaker notes
Task ID: impress_ndo_025
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_025'
OUTPUT = f'{WORKDIR}/Meeting.pptx'

# Slide content and notes
SLIDES = [
    {
        "layout": 0,  # Title Slide
        "title": "Q2 Strategy Meeting",
        "subtitle": "Marketing Department — April 2025",
        "notes": "Welcome everyone to the Q2 strategy meeting. Please ensure all phones are on silent."
    },
    {
        "layout": 1,  # Title + Content
        "title": "Agenda",
        "body": [
            "Q1 Performance Review",
            "Q2 Goals and Objectives",
            "Team Assignments and Roles",
            "Budget Allocation Updates",
            "Next Steps and Timeline",
        ],
        "notes": "Walk through each agenda item briefly. Allocate approximately 10 minutes per topic."
    },
    {
        "layout": 1,
        "title": "Q1 Performance Review",
        "body": [
            "Lead generation up 15% quarter-over-quarter",
            "Social media engagement increased by 22%",
            "Email open rate improved to 34.7%",
            "Website traffic reached 1.2M monthly visits",
            "Customer acquisition cost reduced by $12 per lead",
        ],
        "notes": "Highlight the 15% increase in lead generation. Mention the successful social media campaign."
    },
    {
        "layout": 1,
        "title": "Q2 Goals and Objectives",
        "body": [
            "Launch new product line by May 15th",
            "Increase brand awareness by 20%",
            "Expand into three new regional markets",
            "Achieve $2.4M in quarterly revenue",
            "Reduce customer churn rate to below 5%",
        ],
        "notes": "Emphasize the importance of the product launch timeline. Budget approval is pending from finance."
    },
    {
        "layout": 1,
        "title": "Team Assignments",
        "body": [
            "Sarah Chen — Content Strategy Lead",
            "Marcus Johnson — Analytics and Reporting",
            "Elena Rodriguez — Social Media Campaigns",
            "David Park — Partnership Development",
            "Aisha Thompson — Customer Experience",
        ],
        "notes": "Sarah will lead the content strategy. Marcus handles analytics. Confirm resource allocation with HR."
    },
    {
        "layout": 1,
        "title": "Next Steps and Timeline",
        "body": [
            "April 15 — Follow-up meeting with full team",
            "April 22 — Budget finalization deadline",
            "May 1 — Campaign assets ready for review",
            "May 15 — Product launch date",
            "June 30 — Q2 performance checkpoint",
        ],
        "notes": "Follow-up meeting scheduled for April 15th. All deliverables due by end of March."
    },
]


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    for i, slide_data in enumerate(SLIDES):
        layout_idx = slide_data["layout"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Set title
        slide.shapes.title.text = slide_data["title"]

        # Set subtitle (layout 0) or body content (layout 1)
        if layout_idx == 0 and "subtitle" in slide_data:
            slide.placeholders[1].text = slide_data["subtitle"]
        elif "body" in slide_data:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for j, item in enumerate(slide_data["body"]):
                if j == 0:
                    tf.paragraphs[0].text = item
                else:
                    p = tf.add_paragraph()
                    p.text = item
                    p.level = 0

        # Add speaker notes
        slide.notes_slide.notes_text_frame.text = slide_data["notes"]

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Ensure Desktop directory exists
    os.makedirs(f'{WORKDIR}/Desktop', exist_ok=True)

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
