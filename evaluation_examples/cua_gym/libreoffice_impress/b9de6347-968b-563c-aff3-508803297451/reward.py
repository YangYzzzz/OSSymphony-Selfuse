"""
Reward Script: Export presentation notes as separate text file
Task ID: impress_ndo_025
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): All 6 "Slide X:" headers present in speaker_notes.txt
  Component 2 (0.4): Notes content matches presentation notes per slide (partial credit)
  Component 3 (0.3): Correct format — sections separated by blank lines, proper structure
"""

import os
import re

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_025'
NOTES_FILE = os.path.join(WORKDIR, 'Desktop', 'speaker_notes.txt')
PPTX_FILE = os.path.join(WORKDIR, 'Meeting.pptx')


def get_presentation_notes(pptx_path):
    """Extract notes from the presentation file."""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    notes = {}
    for i, slide in enumerate(prs.slides, 1):
        try:
            text = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            text = ""
        notes[i] = text
    return notes, len(prs.slides)


def verify_task():
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: speaker_notes.txt must exist
    if not os.path.exists(NOTES_FILE):
        print(f"CRITICAL: File not found: {NOTES_FILE}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must exist (to get ground truth notes)
    if not os.path.exists(PPTX_FILE):
        print(f"CRITICAL: Presentation not found: {PPTX_FILE}")
        print("REWARD: 0.0")
        return 0.0

    try:
        with open(NOTES_FILE, 'r') as f:
            content = f.read()
    except Exception as e:
        print(f"CRITICAL: Cannot read {NOTES_FILE}: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        pptx_notes, num_slides = get_presentation_notes(PPTX_FILE)
    except Exception as e:
        print(f"CRITICAL: Cannot read presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Presentation has {num_slides} slides")
    print(f"INFO: Notes file length: {len(content)} chars")

    # Component 1: All 6 "Slide X:" headers present (0.3 points)
    try:
        headers_found = 0
        for i in range(1, num_slides + 1):
            pattern = rf'Slide\s+{i}\s*:'
            if re.search(pattern, content):
                headers_found += 1
            else:
                print(f"FAIL: Component 1 — 'Slide {i}:' header not found")

        if headers_found == num_slides:
            print(f"PASS: Component 1 — All {num_slides} slide headers found (0.3 pts)")
            total_score += 0.3
        elif headers_found > 0:
            partial = round(0.3 * (headers_found / num_slides), 2)
            print(f"PARTIAL: Component 1 — {headers_found}/{num_slides} headers found ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — No slide headers found")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Notes content matches for each slide (0.4 points)
    # Parse the notes file into sections by slide header
    try:
        # Split content into slide sections
        sections = re.split(r'(?=Slide\s+\d+\s*:)', content.strip())
        sections = [s.strip() for s in sections if s.strip()]

        parsed_notes = {}
        for section in sections:
            match = re.match(r'Slide\s+(\d+)\s*:\s*(.*)', section, re.DOTALL)
            if match:
                slide_num = int(match.group(1))
                note_text = match.group(2).strip()
                parsed_notes[slide_num] = note_text

        slides_correct = 0
        per_slide_score = 0.4 / num_slides

        for i in range(1, num_slides + 1):
            expected = pptx_notes.get(i, "").strip()
            actual = parsed_notes.get(i, "").strip()

            if not expected and not actual:
                slides_correct += 1
                print(f"PASS: Component 2 — Slide {i} notes match (both empty)")
            elif expected and actual and expected in actual:
                slides_correct += 1
                print(f"PASS: Component 2 — Slide {i} notes match")
            elif expected and actual and actual in expected:
                slides_correct += 1
                print(f"PASS: Component 2 — Slide {i} notes match (subset)")
            else:
                print(f"FAIL: Component 2 — Slide {i} notes mismatch")
                if expected:
                    print(f"  Expected: {expected[:80]}...")
                if actual:
                    print(f"  Actual:   {actual[:80]}...")
                else:
                    print(f"  Actual:   (missing)")

        if slides_correct > 0:
            earned = round(per_slide_score * slides_correct, 2)
            print(f"{'PASS' if slides_correct == num_slides else 'PARTIAL'}: Component 2 — {slides_correct}/{num_slides} slides correct ({earned} pts)")
            total_score += earned
        else:
            print(f"FAIL: Component 2 — No slide notes matched")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Format correctness (0.3 points)
    # Check that sections are separated by blank lines and follow "Slide X:\n<notes>" format
    try:
        format_checks_passed = 0
        total_format_checks = 3

        # Check 3a: Each "Slide X:" is followed by notes text on next line(s)
        slide_header_pattern = re.compile(r'Slide\s+\d+\s*:\n.+', re.MULTILINE)
        header_with_notes = slide_header_pattern.findall(content)
        if len(header_with_notes) >= num_slides:
            format_checks_passed += 1
            print(f"PASS: Component 3a — Slide headers followed by notes text")
        else:
            print(f"FAIL: Component 3a — Expected {num_slides} headers with notes, found {len(header_with_notes)}")

        # Check 3b: Sections separated by blank lines (double newline between sections)
        # Between each Slide X section, there should be a blank line
        section_gaps = re.findall(r'\n\n+Slide\s+\d+\s*:', content)
        if len(section_gaps) >= num_slides - 1:
            format_checks_passed += 1
            print(f"PASS: Component 3b — Sections separated by blank lines")
        else:
            print(f"FAIL: Component 3b — Expected {num_slides - 1} section separators, found {len(section_gaps)}")

        # Check 3c: Slides appear in order (1, 2, 3, 4, 5, 6)
        slide_numbers = [int(m) for m in re.findall(r'Slide\s+(\d+)\s*:', content)]
        expected_order = list(range(1, num_slides + 1))
        if slide_numbers == expected_order:
            format_checks_passed += 1
            print(f"PASS: Component 3c — Slides in correct order (1-{num_slides})")
        else:
            print(f"FAIL: Component 3c — Slide order: {slide_numbers}, expected: {expected_order}")

        if format_checks_passed > 0:
            earned = round(0.3 * (format_checks_passed / total_format_checks), 2)
            print(f"{'PASS' if format_checks_passed == total_format_checks else 'PARTIAL'}: Component 3 — {format_checks_passed}/{total_format_checks} format checks passed ({earned} pts)")
            total_score += earned
        else:
            print(f"FAIL: Component 3 — All format checks failed")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
verify_task()
