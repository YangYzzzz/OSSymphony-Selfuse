"""
Reward Script: FitTrack 10-slide product pitch verification
Task ID: impress_wf_027
Domain: libreoffice_impress
Scoring:
  C1 (0.15) - 10 slides present
  C2 (0.10) - Slide 1 title with FitTrack + tagline
  C3 (0.10) - Slide 4 has 4 icon placeholder shapes (feature blocks)
  C4 (0.10) - Slide 5 has 3 tall narrow rectangles (phone mockups)
  C5 (0.10) - Slide 6 has 3 quote boxes with names
  C6 (0.15) - Slide 7 has a doughnut chart
  C7 (0.10) - Slide 8 has a table
  C8 (0.10) - Slide 9 has horizontal arrow + 4 marker shapes
  C9 (0.10) - Green accent #4CAF50 used on multiple slides
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_027'
FILE_NAME = 'FitTrack_Pitch.pptx'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)
    num_slides = len(slides)

    # Component 1: Exactly 10 slides (0.15 points)
    try:
        if num_slides == 10:
            print(f"PASS: Component 1 — 10 slides found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — Expected 10 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Guard: need at least 10 slides for remaining checks
    if num_slides < 10:
        print(f"\nInsufficient slides for remaining checks.")
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: Slide 1 has 'FitTrack' and 'Your Personal Fitness Journey' (0.10 pts)
    try:
        slide1_text = ""
        for shape in slides[0].shapes:
            if shape.has_text_frame:
                slide1_text += " " + shape.text_frame.text
        has_fittrack = "FitTrack" in slide1_text or "fittrack" in slide1_text.lower()
        has_tagline = "your personal fitness journey" in slide1_text.lower()
        if has_fittrack and has_tagline:
            print(f"PASS: Component 2 — Slide 1 has FitTrack title + tagline (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 — FitTrack={has_fittrack}, tagline={has_tagline}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 4 has 4 icon placeholder shapes (feature blocks) (0.10 pts)
    # The task says 4 features with icon placeholders; golden has 4 rounded rectangle auto_shapes
    try:
        slide4 = slides[3]
        icon_placeholders = 0
        for shape in slide4.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # Count rounded/regular rectangles that are roughly square (icon-like)
                # Exclude the top accent bar (very wide, very short)
                w = shape.width
                h = shape.height
                # Icon placeholders are roughly square; accent bars are very wide
                if w > 0 and h > 0:
                    ratio = max(w, h) / min(w, h)
                    if ratio < 3:  # roughly square-ish, not a thin bar
                        icon_placeholders += 1
        if icon_placeholders >= 4:
            print(f"PASS: Component 3 — Slide 4 has {icon_placeholders} icon placeholders (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 3 — Expected 4 icon placeholders, found {icon_placeholders}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 5 has 3 tall narrow rectangles (phone mockups) (0.10 pts)
    # Phone mockups: height >> width (tall and narrow)
    try:
        slide5 = slides[4]
        phone_shapes = 0
        for shape in slide5.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                w = shape.width
                h = shape.height
                if w > 0 and h > 0 and h > w * 1.5:
                    # Tall and narrow — phone mockup
                    phone_shapes += 1
        if phone_shapes >= 3:
            print(f"PASS: Component 4 — Slide 5 has {phone_shapes} phone mockup shapes (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 4 — Expected 3 phone mockups, found {phone_shapes}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slide 6 has 3 quote boxes with attribution names (0.10 pts)
    try:
        slide6 = slides[5]
        # Count text shapes that contain a dash/em-dash followed by a name (attribution)
        name_attributions = 0
        for shape in slide6.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                # Attribution lines start with dash/em-dash/hyphen
                if text and (text.startswith('—') or text.startswith('-') or text.startswith('\u2014')):
                    name_attributions += 1
        if name_attributions >= 3:
            print(f"PASS: Component 5 — Slide 6 has {name_attributions} quote attributions (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 5 — Expected 3 quote attributions, found {name_attributions}")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Slide 7 has a doughnut chart (0.15 pts)
    try:
        slide7 = slides[6]
        has_doughnut = False
        for shape in slide7.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                # DOUGHNUT chart type value is -4120
                ct = shape.chart.chart_type
                if ct is not None and int(ct) == -4120:
                    has_doughnut = True
                    break
        if has_doughnut:
            print(f"PASS: Component 6 — Slide 7 has doughnut chart (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 6 — Slide 7 missing doughnut chart")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Slide 8 has a table (0.10 pts)
    try:
        slide8 = slides[7]
        has_table = False
        table_rows = 0
        table_cols = 0
        for shape in slide8.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                has_table = True
                table_rows = len(shape.table.rows)
                table_cols = len(shape.table.columns)
                break
        if has_table and table_rows >= 2 and table_cols >= 2:
            print(f"PASS: Component 7 — Slide 8 has table ({table_rows}x{table_cols}) (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 7 — Slide 8 table: present={has_table}, rows={table_rows}, cols={table_cols}")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    # Component 8: Slide 9 has horizontal arrow + 4 marker shapes (0.10 pts)
    try:
        slide9 = slides[8]
        has_arrow = False
        marker_count = 0
        for shape in slide9.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                name_lower = shape.name.lower()
                w = shape.width
                h = shape.height
                # Arrow: much wider than tall
                if ('arrow' in name_lower) and w > h * 2:
                    has_arrow = True
                # Markers: oval/circle shapes (roughly equal width and height)
                elif w > 0 and h > 0:
                    ratio = max(w, h) / min(w, h)
                    if ratio < 2 and w < slide9.shapes[0].width * 0.5:
                        # Small-ish roughly-square shapes = markers
                        marker_count += 1
        if has_arrow and marker_count >= 4:
            print(f"PASS: Component 8 — Slide 9 has arrow + {marker_count} markers (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 8 — arrow={has_arrow}, markers={marker_count}")
    except Exception as e:
        print(f"ERROR: Component 8 — {e}")

    # Component 9: Green accent #4CAF50 used on at least 3 slides (0.10 pts)
    try:
        green_hex = '4CAF50'
        slides_with_green = set()
        for i, slide in enumerate(slides):
            found_green = False
            for shape in slide.shapes:
                # Check text color
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                if run.font.color.type is not None:
                                    if green_hex in str(run.font.color.rgb):
                                        found_green = True
                            except:
                                pass
                # Check shape fill
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    try:
                        f = shape.fill
                        if f.type is not None and f.type == 1:
                            if green_hex in str(f.fore_color.rgb):
                                found_green = True
                    except:
                        pass
                if found_green:
                    slides_with_green.add(i + 1)
                    break

        num_green_slides = len(slides_with_green)
        if num_green_slides >= 3:
            print(f"PASS: Component 9 — Green #4CAF50 found on {num_green_slides} slides: {sorted(slides_with_green)} (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 9 — Green #4CAF50 found on only {num_green_slides} slides")
    except Exception as e:
        print(f"ERROR: Component 9 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point: find the file on Desktop
file_path = os.path.join(WORKDIR, 'Desktop', FILE_NAME)
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
