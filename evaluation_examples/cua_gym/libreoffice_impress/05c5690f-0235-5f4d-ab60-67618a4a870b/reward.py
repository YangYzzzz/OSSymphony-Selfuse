"""
Reward Script: Set subtitle text on slide 1 to italic 18pt Georgia, color #757575, 1.5 line spacing
Task ID: impress_tct_072
Domain: libreoffice_impress
Scoring:
  Component 1: Font is italic (0.2)
  Component 2: Font size is 18pt (0.2)
  Component 3: Font name is Georgia (0.2)
  Component 4: Font color is #757575 (0.2)
  Component 5: Line spacing is 1.5 (0.2)
"""

import os
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_072'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 1 slide
    if len(prs.slides) < 1:
        print("FAIL: Presentation has no slides")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]

    # Find the subtitle placeholder (named "Subtitle 2" or placeholder index 1)
    subtitle_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and "subtitle" in shape.name.lower():
            subtitle_shape = shape
            break

    if subtitle_shape is None:
        # Fallback: try placeholder index 1
        try:
            subtitle_shape = slide.placeholders[1]
        except (KeyError, IndexError):
            pass

    if subtitle_shape is None or not subtitle_shape.has_text_frame:
        print("FAIL: No subtitle shape found on slide 1")
        print("REWARD: 0.0")
        return 0.0

    tf = subtitle_shape.text_frame

    # Precondition: text content is preserved
    full_text = tf.text.strip()
    if "Transforming the Future of Finance" not in full_text:
        print(f"FAIL: Subtitle text changed. Expected 'Transforming the Future of Finance', found: {repr(full_text)}")
        print("REWARD: 0.0")
        return 0.0

    # Get the first paragraph and its runs
    para = tf.paragraphs[0]
    runs = [r for r in para.runs if (r.text or "").strip()]

    if not runs:
        print("FAIL: No non-empty runs in subtitle paragraph")
        print("REWARD: 0.0")
        return 0.0

    # We check properties on the first run (all runs should have the same formatting)
    run = runs[0]

    # Component 1: Font is italic (0.2 points)
    try:
        is_italic = run.font.italic
        # None means inherited (not explicitly set), treat as not italic
        if is_italic is True:
            print(f"PASS: Component 1 - Font is italic (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 1 - Expected italic=True, found italic={is_italic}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Font size is 18pt (228600 EMU) (0.2 points)
    try:
        font_size = run.font.size
        expected_size = Pt(18)  # 228600 EMU
        if font_size is not None and font_size == expected_size:
            print(f"PASS: Component 2 - Font size is 18pt ({font_size} EMU) (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 2 - Expected size={expected_size} (18pt), found size={font_size}")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Font name is Georgia (0.2 points)
    try:
        font_name = run.font.name
        if font_name is not None and font_name.lower() == "georgia":
            print(f"PASS: Component 3 - Font name is Georgia (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 3 - Expected font name='Georgia', found '{font_name}'")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Font color is #757575 (0.2 points)
    try:
        color_type = run.font.color.type
        if color_type is not None:
            color_rgb = str(run.font.color.rgb).upper()
            if color_rgb == "757575":
                print(f"PASS: Component 4 - Font color is #757575 (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 4 - Expected color=#757575, found #{color_rgb}")
        else:
            print(f"FAIL: Component 4 - Font color type is None (inherited/not set)")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Line spacing is 1.5 (150000 pct value in XML) (0.2 points)
    try:
        pPr = para._p.find(qn('a:pPr'))
        lnSpc = pPr.find(qn('a:lnSpc')) if pPr is not None else None
        if lnSpc is not None:
            spc_pct = lnSpc.find(qn('a:spcPct'))
            if spc_pct is not None:
                val = int(spc_pct.get('val', '0'))
                if val == 150000:
                    print(f"PASS: Component 5 - Line spacing is 1.5 (150000 pct) (0.2 pts)")
                    total_score += 0.2
                else:
                    print(f"FAIL: Component 5 - Expected line spacing pct=150000, found {val}")
            else:
                # Check for spcPts (points-based spacing)
                spc_pts = lnSpc.find(qn('a:spcPts'))
                if spc_pts is not None:
                    pts_val = int(spc_pts.get('val', '0'))
                    # 1.5 line spacing at 18pt = 27pt = 2700 hundredths
                    print(f"FAIL: Component 5 - Line spacing is in points ({pts_val}), expected percentage 150000")
                else:
                    print(f"FAIL: Component 5 - Line spacing element found but no spcPct or spcPts child")
        else:
            print(f"FAIL: Component 5 - No line spacing element found (default single spacing)")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
