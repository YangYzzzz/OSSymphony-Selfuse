"""
Initial Setup: Investor Pitch presentation with 10 slides.
Slide 1 subtitle is 14pt regular Liberation Sans, black, single line spacing.
Task ID: impress_tct_072
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_072'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title Slide ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    # Title
    title_shape = slide1.shapes.title
    title_shape.text = ""
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "FinEdge Capital"
    run.font.name = "Liberation Sans"
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # Subtitle — 14pt regular Liberation Sans, black, single line spacing
    subtitle_shape = slide1.placeholders[1]
    subtitle_shape.text = ""
    stf = subtitle_shape.text_frame
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    # Ensure single line spacing (default, but set explicitly for clarity)
    from pptx.oxml.ns import qn
    pPr = sp._p.get_or_add_pPr()
    lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
    spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': '100000'})  # 100% = single
    lnSpc.append(spcPct)
    pPr.append(lnSpc)

    srun = sp.add_run()
    srun.text = "Transforming the Future of Finance"
    srun.font.name = "Liberation Sans"
    srun.font.size = Pt(14)
    srun.font.bold = False
    srun.font.italic = False
    srun.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # ── Slide 2: Mission & Vision ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Our Mission & Vision"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Mission: Democratize sophisticated financial tools for mid-market enterprises"
    p2 = body2.add_paragraph()
    p2.text = "Vision: Become the leading AI-powered financial analytics platform by 2028"
    p3 = body2.add_paragraph()
    p3.text = "We bridge the gap between institutional-grade analytics and accessible SaaS solutions"

    # ── Slide 3: Market Opportunity ──
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Market Opportunity"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Total Addressable Market: $47.2B (2025)"
    for line in [
        "Serviceable Addressable Market: $12.8B",
        "Mid-market segment growing at 23% CAGR",
        "87% of mid-market CFOs cite analytics gaps as top pain point",
        "Legacy solutions cost 5-10x more with 3x longer implementation",
    ]:
        p = body3.add_paragraph()
        p.text = line

    # ── Slide 4: Product Overview ──
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Product Overview"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "FinEdge Analytics Suite"
    for line in [
        "Real-time cash flow forecasting with 94% accuracy",
        "Automated regulatory compliance reporting (SOX, IFRS)",
        "AI-driven anomaly detection across 200+ data sources",
        "Natural language query interface for non-technical users",
        "Seamless integration with SAP, Oracle, and QuickBooks",
    ]:
        p = body4.add_paragraph()
        p.text = line

    # ── Slide 5: Business Model ──
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Business Model"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "SaaS subscription tiers:"
    for line in [
        "Starter: $2,500/month — up to 50 users, core analytics",
        "Professional: $7,500/month — unlimited users, AI insights",
        "Enterprise: $18,000/month — custom integrations, dedicated support",
        "Average contract value: $112,000 ARR",
        "Net revenue retention: 135%",
    ]:
        p = body5.add_paragraph()
        p.text = line

    # ── Slide 6: Traction & Metrics ──
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Traction & Key Metrics"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Q4 2025 Highlights"
    for line in [
        "ARR: $8.4M (up 180% YoY)",
        "Active customers: 127 mid-market enterprises",
        "Monthly active users: 14,200",
        "Customer acquisition cost: $18,500",
        "LTV/CAC ratio: 6.1x",
        "Gross margin: 82%",
    ]:
        p = body6.add_paragraph()
        p.text = line

    # ── Slide 7: Competitive Landscape ──
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Competitive Landscape"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "Key differentiators vs. incumbents:"
    for line in [
        "Bloomberg Terminal: Enterprise-only, $24K/user/year — we serve mid-market at 1/10th cost",
        "Anaplan: 6-month implementation — our average onboarding is 3 weeks",
        "Tableau: Visualization-only — we provide predictive analytics + compliance",
        "Adaptive Insights: Limited AI — our ML models trained on 2.3B financial transactions",
    ]:
        p = body7.add_paragraph()
        p.text = line

    # ── Slide 8: Go-to-Market Strategy ──
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Go-to-Market Strategy"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Three-pronged approach:"
    for line in [
        "Channel 1: Direct sales team (12 AEs targeting $100M-$2B revenue companies)",
        "Channel 2: Strategic partnerships with Big 4 accounting firms",
        "Channel 3: Product-led growth via free compliance audit tool",
        "Target: 300 enterprise customers by end of 2026",
    ]:
        p = body8.add_paragraph()
        p.text = line

    # ── Slide 9: Team ──
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Leadership Team"
    body9 = slide9.placeholders[1].text_frame
    body9.text = "Experienced operators from finance and technology:"
    for line in [
        "Dr. Elena Vasquez, CEO — Former VP of Product at Stripe, Stanford MBA",
        "James Okafor, CTO — Ex-Google Staff Engineer, ML at Scale expertise",
        "Rachel Kim, CFO — Former Controller at Palantir, CPA",
        "David Moretti, VP Sales — Built $50M ARR book at Salesforce",
        "Advisory Board: Partners from Sequoia, a16z, and Goldman Sachs",
    ]:
        p = body9.add_paragraph()
        p.text = line

    # ── Slide 10: The Ask ──
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "Series B: $35M Raise"
    body10 = slide10.placeholders[1].text_frame
    body10.text = "Use of proceeds:"
    for line in [
        "40% — Engineering & Product (expand ML team, APAC data centers)",
        "30% — Sales & Marketing (grow GTM team from 15 to 45)",
        "20% — Strategic acquisitions (compliance data providers)",
        "10% — Operations & G&A",
        "Target close: Q2 2026",
        "Expected runway: 24 months to profitability",
    ]:
        p = body10.add_paragraph()
        p.text = line

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
