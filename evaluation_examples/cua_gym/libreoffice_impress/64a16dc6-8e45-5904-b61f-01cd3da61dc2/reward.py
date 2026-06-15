"""
FINAL REWARD SCRIPT - SUCCESS
Task: The bullet list on slide 168 feels squished together—there’s no breathing room after each point. How can I set the paragraph spacing for those bullets to 0 pt before and 8 pt after in LibreOffice Impress?
Generated: 2025-09-10 15:27:19
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.exc import PackageNotFoundError


def verify_paragraph_spacing(
    file_path: str,
    slide_index: int = 167,  # zero-based index; 167 means slide 168
    expected_before_pt: float = 0.0,
    expected_after_pt: float = 8.0,
    tolerance: float = 0.5,
) -> float:
    """Verify that every bullet paragraph on a specific slide has the
    required paragraph spacing (0 pt before, 8 pt after by default).

    Returns a progressive score between 0.0 and 1.0.
    """
    print(f"Loading presentation: {file_path}")

    # --- prerequisite check: file must exist and load successfully ---
    if not os.path.exists(file_path):
        print("✗ File not found.")
        return 0.0

    try:
        prs = Presentation(file_path)
    except PackageNotFoundError as e:
        print(f"✗ Unable to open presentation: {e}")
        return 0.0

    total_slides = len(prs.slides)
    print(f"✓ Presentation loaded with {total_slides} slides")

    # --- verify the requested slide actually exists ---
    if slide_index >= total_slides:
        print(f"✗ Slide {slide_index + 1} not found in presentation.")
        return 0.0

    slide = prs.slides[slide_index]

    # --- collect bullet-style paragraphs (exclude title placeholder) ---
    bullet_paragraphs = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        # skip title placeholders – they are not part of the bullet list
        if shape.is_placeholder and shape.placeholder_format.type in {
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
        }:
            continue

        for p in shape.text_frame.paragraphs:
            if not p.text.strip():
                continue  # ignore empty lines
            bullet_paragraphs.append(p)

    print(
        f"Found {len(bullet_paragraphs)} bullet paragraphs on slide {slide_index + 1}."
    )

    # ---------------- progressive scoring ----------------
    total_score = 0.0

    # 1) Must have at least one bullet paragraph (0.4 points)
    if bullet_paragraphs:
        print("✓ Bullet paragraphs detected (0.4 points)")
        total_score += 0.4
    else:
        print("✗ No bullet paragraphs found – 0 points")
        return 0.0  # cannot evaluate spacing without content

    # 2) Verify space_before (0.3 points if all match)
    before_ok = True
    for p in bullet_paragraphs:
        before_len = p.space_before
        before_pt = before_len.pt if before_len else 0.0
        if abs(before_pt - expected_before_pt) > tolerance:
            before_ok = False
            print(
                f"  ✗ Paragraph '{p.text[:30]}…' has space_before {before_pt} pt (expected {expected_before_pt})"
            )
        else:
            print(
                f"  ✓ Paragraph '{p.text[:30]}…' has correct space_before {before_pt} pt"
            )

    if before_ok:
        print("✓ All bullet paragraphs have correct space_before (0.3 points)")
        total_score += 0.3
    else:
        print("✗ Some paragraphs have incorrect space_before – 0 points for this section")

    # 3) Verify space_after (0.3 points if all match)
    after_ok = True
    for p in bullet_paragraphs:
        after_len = p.space_after
        after_pt = after_len.pt if after_len else 0.0
        if abs(after_pt - expected_after_pt) > tolerance:
            after_ok = False
            print(
                f"  ✗ Paragraph '{p.text[:30]}…' has space_after {after_pt} pt (expected {expected_after_pt})"
            )
        else:
            print(
                f"  ✓ Paragraph '{p.text[:30]}…' has correct space_after {after_pt} pt"
            )

    if after_ok:
        print("✓ All bullet paragraphs have correct space_after (0.3 points)")
        total_score += 0.3
    else:
        print("✗ Some paragraphs have incorrect space_after – 0 points for this section")

    # --- final score (capped at 1.0) ---
    final_score = min(total_score, 1.0)
    print(f"Total score: {final_score}")
    return final_score


if __name__ == "__main__":
    FILE_PATH = "/home/user/the_bullet_list_on_slide_168_feels_squished_togethertheres_no_breathing_room_after_each_point_how_ca_golden.pptx"
    reward = verify_paragraph_spacing(FILE_PATH)
    print(f"REWARD: {reward}")

