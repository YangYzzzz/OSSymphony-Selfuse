"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 3 of my LibreOffice Impress file there’s a text box whose title line literally says "Personal Info" in 18 pt Arial. I need that entire box gone—content, border, everything—so it doesn’t show up when the deck hits slide 100 during the presentation.
Generated: 2025-09-10 15:32:43
Status: success
Model: azure-o3
Total Steps: 1
"""

from pptx import Presentation
import os

# ----------------- Helper Functions -----------------

def extract_runs_from_shape(shape):
    """Recursively extract all (text, font_size_pt, font_name) tuples from a shape.

    Handles grouped shapes by drilling down into sub-shapes and only returns
    details for shapes that actually contain a text-frame.
    """
    runs = []

    # If the shape itself has no text frame, but is a group, recurse
    if not getattr(shape, "has_text_frame", False):
        if shape.shape_type == 6:  # 6 == GROUP
            for sub_shape in shape.shapes:
                runs.extend(extract_runs_from_shape(sub_shape))
        return runs

    # Shape HAS a text frame – iterate through its paragraphs and runs
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = run.text or ""
            size_pt = run.font.size.pt if run.font.size is not None else None
            font_name = run.font.name
            runs.append((text, size_pt, font_name))
    return runs

# ----------------- Main Verification -----------------

def verify_task(file_path):
    """Verify that the text box titled *Personal Info* is completely removed from slide 3.

    Scoring (progressive):
    0.7  – No text containing the phrase "Personal Info" exists on slide 3
    0.3  – No run on slide 3 is formatted with 18 pt Arial (original style)
    1.0  – Both conditions satisfied
    """
    score = 0.0
    max_score = 1.0

    print(f"Verifying presentation: {file_path}")

    # ---------- Basic file checks (NO points awarded) ----------
    if not os.path.exists(file_path):
        print("✗ File not found – cannot verify task")
        return 0.0
    if not file_path.lower().endswith(".pptx"):
        print("✗ Unsupported file type – expected .pptx")
        return 0.0

    # ---------- Attempt to load the presentation ----------
    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded PPTX with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        return 0.0

    # ---------- Ensure slide 3 exists ----------
    if len(prs.slides) < 3:
        print("✗ Presentation has fewer than 3 slides – cannot verify")
        return 0.0

    slide3 = prs.slides[2]  # zero-based index

    # ---------- Gather ALL text runs from slide 3 ----------
    all_runs = []
    for shp in slide3.shapes:
        all_runs.extend(extract_runs_from_shape(shp))

    # ---------- Condition 1: Phrase "Personal Info" removed ----------
    phrase_present = any("personal info" in (txt.lower()) for txt, _, _ in all_runs)
    if phrase_present:
        print("✗ Found forbidden phrase 'Personal Info' on slide 3 (0 points)")
    else:
        print("✓ Phrase 'Personal Info' absent on slide 3 (0.7 points)")
        score += 0.7

    # ---------- Condition 2: Original 18 pt Arial styling removed ----------
    style_present = False
    for txt, size_pt, font_name in all_runs:
        if size_pt is None:
            continue
        # Tolerate small rounding differences (<0.5 pt)
        if abs(size_pt - 18) < 0.5 and font_name and font_name.lower() == "arial":
            style_present = True
            print(f"   Found 18 pt Arial text run: '{txt.strip()}'")
            break

    if style_present:
        print("✗ 18 pt Arial formatting still present on slide 3 (0 points)")
    else:
        print("✓ 18 pt Arial formatting absent on slide 3 (0.3 points)")
        score += 0.3

    # ---------- Final score ----------
    final_score = round(min(score, max_score), 2)
    print(f"Total score: {final_score}")
    return final_score

# ----------------- Execute Verification -----------------

if __name__ == "__main__":
    FILE_PATH = "/home/user/on_slide_3_of_my_libreoffice_impress_file_theres_a_text_box_whose_title_line_literally_says_personal_golden.pptx"
    reward = verify_task(FILE_PATH)
    print(f"REWARD: {reward}")
