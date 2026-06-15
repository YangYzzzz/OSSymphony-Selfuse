"""
Initial Setup: Create a 12-slide tech demo presentation with no transitions
Task ID: impress_tm_035
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_035'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    """Add a title slide (layout 0)."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points):
    """Add a title+content slide (layout 1)."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(18)
    return slide


def add_blank_slide_with_text(prs, title_text, body_text):
    """Add a blank slide with a title textbox and body textbox."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
    # Body
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = body_text
    for run in p2.runs:
        run.font.size = Pt(16)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(prs,
                    "CloudMatrix Tech Demo",
                    "Next-Generation Infrastructure Platform\nQ2 2025 Technical Review")

    # Slide 2: Agenda
    add_content_slide(prs, "Agenda", [
        "Platform Overview & Vision",
        "Core Technology Stack",
        "Performance Benchmarks",
        "Security Framework",
        "Deployment Pipeline",
        "Live Demo Environment",
        "Architecture Overview",
        "Scalability Metrics",
        "Roadmap & Next Steps",
    ])

    # Slide 3: Platform Overview
    add_content_slide(prs, "Platform Overview", [
        "Multi-cloud orchestration across AWS, Azure, and GCP",
        "Real-time data processing with sub-10ms latency",
        "Auto-scaling from 100 to 50,000 concurrent connections",
        "99.99% uptime SLA with active-active failover",
        "Built on Kubernetes 1.28 with custom operators",
    ])

    # Slide 4: Core Technology Stack
    add_content_slide(prs, "Core Technology Stack", [
        "Backend: Go 1.22 microservices with gRPC",
        "Frontend: React 18 + TypeScript with SSR",
        "Database: PostgreSQL 16 + Redis 7 cluster",
        "Message Queue: Apache Kafka 3.7 with exactly-once semantics",
        "Observability: OpenTelemetry + Grafana + Prometheus",
    ])

    # Slide 5: Performance Benchmarks
    add_content_slide(prs, "Performance Benchmarks", [
        "API Response Time: P50 = 8ms, P95 = 23ms, P99 = 45ms",
        "Throughput: 180,000 requests/second per node",
        "Data Ingestion: 2.4 million events/second sustained",
        "Cold Start: < 350ms for new container instances",
        "Memory Efficiency: 40% reduction vs. previous generation",
    ])

    # Slide 6: Security Framework
    add_content_slide(prs, "Security Framework", [
        "Zero-trust network architecture with mTLS everywhere",
        "FIDO2/WebAuthn for developer authentication",
        "Hardware Security Modules for key management",
        "SOC 2 Type II and ISO 27001 certified",
        "Automated vulnerability scanning in CI/CD pipeline",
    ])

    # Slide 7: Deployment Pipeline
    add_content_slide(prs, "Deployment Pipeline", [
        "GitOps workflow with ArgoCD and Flux",
        "Canary deployments with automatic rollback",
        "Blue-green deployments for zero-downtime releases",
        "Average deployment time: 4 minutes 12 seconds",
        "300+ deployments per week across all services",
    ])

    # Slide 8: Live Demo Environment
    add_blank_slide_with_text(prs,
                              "Live Demo Environment",
                              "The demo environment runs on a dedicated Kubernetes cluster with "
                              "3 control-plane nodes and 12 worker nodes. Each worker node has "
                              "32 vCPUs and 128 GB RAM, providing sufficient capacity for "
                              "realistic load testing. The environment mirrors production "
                              "configuration with TLS termination at the ingress controller "
                              "and service mesh communication via Istio 1.20.")

    # Slide 9: Architecture Overview (NO TRANSITION - this is the target slide)
    add_blank_slide_with_text(prs,
                              "Architecture Overview",
                              "The CloudMatrix platform follows a hexagonal architecture pattern "
                              "with clear boundaries between domain logic and infrastructure. "
                              "Inbound adapters handle REST, gRPC, and WebSocket protocols while "
                              "outbound adapters manage database connections, message publishing, "
                              "and external API integrations. The core domain layer remains "
                              "transport-agnostic, enabling independent testing and future "
                              "protocol additions without modifying business logic.")

    # Slide 10: Scalability Metrics
    add_content_slide(prs, "Scalability Metrics", [
        "Horizontal scaling: 2 to 200 pods in under 90 seconds",
        "Database read replicas: automatic promotion in < 30s",
        "CDN cache hit ratio: 94.7% average across all regions",
        "Global edge locations: 42 points of presence",
        "Cross-region replication lag: < 150ms P99",
    ])

    # Slide 11: Customer Success Stories
    add_content_slide(prs, "Customer Success Stories", [
        "Meridian Financial: 60% reduction in infrastructure costs",
        "Nova Healthcare: HIPAA-compliant deployment in 3 weeks",
        "Zenith Retail: Black Friday traffic handled with zero incidents",
        "Atlas Logistics: Real-time tracking for 500K daily shipments",
    ])

    # Slide 12: Roadmap & Next Steps
    add_content_slide(prs, "Roadmap & Next Steps", [
        "Q3 2025: AI-powered auto-tuning for resource allocation",
        "Q4 2025: Multi-region active-active database support",
        "Q1 2026: Edge computing integration with 5G networks",
        "Q2 2026: Quantum-safe encryption migration begins",
        "Ongoing: Developer experience improvements and SDK updates",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    subprocess.run(["pkill", "-9", "-f", "soffice"], capture_output=True)
    time.sleep(1)
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=3.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
