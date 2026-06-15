"""
Initial Setup: Create design_showcase.pptx with 10 slides for a UX design showcase.
Task ID: impress_gf5_016
Domain: libreoffice_impress

Slides 3, 5, and 7 use the Blank layout. No custom layouts exist.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_016'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_shape(shape, text, font_size=Pt(18), bold=False, color=None, alignment=None):
    """Helper to set text on a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_bullet_points(text_frame, items, font_size=Pt(14)):
    """Add bullet point items to a text frame."""
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = font_size


def create_initial():
    prs = Presentation()
    # Standard slide dimensions
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Get layout references
    # 0=Title Slide, 1=Title+Content, 5=Blank, 6=Title Only
    layout_title = prs.slide_layouts[0]
    layout_content = prs.slide_layouts[1]
    layout_blank = prs.slide_layouts[6]
    layout_title_only = prs.slide_layouts[5]

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(layout_title)
    slide1.shapes.title.text = "NexaWave Design Showcase"
    slide1.placeholders[1].text = "Product Design & Feature Overview\nQ2 2025 Sprint Review"

    # ---- Slide 2: Agenda (Title+Content) ----
    slide2 = prs.slides.add_slide(layout_content)
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    add_bullet_points(body2, [
        "Dashboard Redesign Overview",
        "New Feature: Smart Filters",
        "Mobile Responsiveness Updates",
        "User Feedback Integration",
        "Analytics Engine Improvements",
        "Accessibility Compliance Progress",
        "Next Steps & Timeline",
    ])

    # ---- Slide 3: BLANK (target for Feature Slide layout) ----
    slide3 = prs.slides.add_slide(layout_blank)
    # Add some placeholder content as text boxes (not layout placeholders)
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Smart Filter Panel"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2D, 0x3A, 0x4A)

    txBox2 = slide3.shapes.add_textbox(Inches(0.5), Inches(2), Inches(8), Inches(4))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    add_bullet_points(tf2, [
        "Contextual filtering based on user behavior patterns",
        "Real-time search suggestions with machine learning",
        "Saved filter presets for repeat workflows",
        "Cross-module filter synchronization",
    ])

    # ---- Slide 4: Title+Content ----
    slide4 = prs.slides.add_slide(layout_content)
    slide4.shapes.title.text = "Dashboard Redesign Metrics"
    body4 = slide4.placeholders[1].text_frame
    add_bullet_points(body4, [
        "Page load time reduced from 3.2s to 1.1s",
        "User engagement up 47% after redesign",
        "Support ticket volume decreased by 32%",
        "Net Promoter Score improved from 34 to 61",
        "Average session duration increased by 2.5 minutes",
    ])

    # ---- Slide 5: BLANK (target for Feature Slide layout) ----
    slide5 = prs.slides.add_slide(layout_blank)
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Mobile Responsive Layout"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2D, 0x3A, 0x4A)

    txBox2 = slide5.shapes.add_textbox(Inches(0.5), Inches(2), Inches(8), Inches(4))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    add_bullet_points(tf2, [
        "Adaptive grid system for tablets and phones",
        "Touch-optimized controls with haptic feedback",
        "Offline mode with local data caching",
        "Progressive image loading for slow connections",
    ])

    # ---- Slide 6: Title+Content ----
    slide6 = prs.slides.add_slide(layout_content)
    slide6.shapes.title.text = "User Feedback Summary"
    body6 = slide6.placeholders[1].text_frame
    add_bullet_points(body6, [
        "892 survey responses collected (78% response rate)",
        "Top request: bulk action support in data tables",
        "Second request: dark mode for extended sessions",
        "Pain point: notification overload during peak hours",
        "Positive: new search ranked 4.6/5 by beta testers",
    ])

    # ---- Slide 7: BLANK (target for Feature Slide layout) ----
    slide7 = prs.slides.add_slide(layout_blank)
    txBox = slide7.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Analytics Engine v2"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2D, 0x3A, 0x4A)

    txBox2 = slide7.shapes.add_textbox(Inches(0.5), Inches(2), Inches(8), Inches(4))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    add_bullet_points(tf2, [
        "Real-time data pipeline with sub-second latency",
        "Custom dashboard widgets with drag-and-drop",
        "Automated anomaly detection and alerting",
        "Export to PDF, Excel, and scheduled email reports",
    ])

    # ---- Slide 8: Title+Content ----
    slide8 = prs.slides.add_slide(layout_content)
    slide8.shapes.title.text = "Accessibility Compliance"
    body8 = slide8.placeholders[1].text_frame
    add_bullet_points(body8, [
        "WCAG 2.1 AA compliance achieved for core flows",
        "Screen reader compatibility verified with NVDA and VoiceOver",
        "Color contrast ratios meet 4.5:1 minimum standard",
        "Keyboard navigation fully operational across all modules",
        "Remaining: AAA compliance for high-priority pages",
    ])

    # ---- Slide 9: Title+Content ----
    slide9 = prs.slides.add_slide(layout_content)
    slide9.shapes.title.text = "Sprint Velocity & Timeline"
    body9 = slide9.placeholders[1].text_frame
    add_bullet_points(body9, [
        "Sprint 14: 87 story points completed (team average: 72)",
        "Feature freeze date: June 15, 2025",
        "QA cycle: June 16 - June 30, 2025",
        "Staged rollout begins: July 7, 2025",
        "Full deployment target: July 21, 2025",
    ])

    # ---- Slide 10: Title Slide (Closing) ----
    slide10 = prs.slides.add_slide(layout_title)
    slide10.shapes.title.text = "Thank You"
    slide10.placeholders[1].text = "Questions & Discussion\ndesign@nexawave.io"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
