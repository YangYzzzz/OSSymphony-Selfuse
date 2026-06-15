"""
Reward Script: Custom 'Feature Slide' layout in Slide Master
Task ID: impress_gf5_016
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): 'Feature Slide' layout exists in the slide master
  Component 2 (0.30): Layout has correct placeholders (title top, picture left half, content right half)
  Component 3 (0.40): Slides 3, 5, 7 use the 'Feature Slide' layout; other slides unaffected
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_016'


def persist_app_state(domain: str):
    """Best-effort save of any open LibreOffice document."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(1.0)
            print("PERSIST: ctrl+s sent for libreoffice_impress")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file has 10 slides
    if len(prs.slides) != 10:
        print(f"PRECONDITION FAIL: Expected 10 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide_width = prs.slide_width

    # -------------------------------------------------------------------
    # Component 1: 'Feature Slide' layout exists in the slide master (0.30)
    # -------------------------------------------------------------------
    feature_layout = None
    try:
        layout_names = []
        for layout in prs.slide_layouts:
            layout_names.append(layout.name)
            if layout.name == 'Feature Slide':
                feature_layout = layout

        if feature_layout is not None:
            print(f"PASS: Component 1 — 'Feature Slide' layout found in slide master (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 1 — 'Feature Slide' layout not found. Available layouts: {layout_names}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # -------------------------------------------------------------------
    # Component 2: Layout has correct placeholders (0.30)
    #   - Title placeholder (type TITLE) at the top
    #   - Picture/image placeholder on the left ~half
    #   - Content/Object placeholder on the right ~half
    # -------------------------------------------------------------------
    if feature_layout is not None:
        try:
            placeholders = list(feature_layout.placeholders)
            # Filter out date/footer/slide number placeholders (idx >= 10)
            content_phs = [ph for ph in placeholders if ph.placeholder_format.idx < 10]

            has_title = False
            has_image_left = False
            has_content_right = False
            sub_score = 0.0

            for ph in content_phs:
                pf = ph.placeholder_format
                ph_type = pf.type
                # Check for title placeholder (type 1 = TITLE)
                # Title should be near the top of the slide
                if ph_type == 1:  # TITLE
                    has_title = True

                # Check for picture/image placeholder (type 18 = PICTURE)
                # Should be on the left half of the slide
                elif ph_type == 18:  # PICTURE
                    right_edge = ph.left + ph.width
                    half_width = slide_width / 2
                    # Image placeholder right edge should be roughly at or before the slide midpoint
                    # Allow some tolerance (within 20% of slide width from midpoint)
                    if right_edge <= half_width * 1.2:
                        has_image_left = True
                    else:
                        print(f"  INFO: Picture PH right_edge={right_edge}, slide_half={half_width} — not on left half")

                # Check for content/object placeholder (type 7 = OBJECT or type 2 = BODY)
                # Should be on the right half of the slide
                elif ph_type in (7, 2):  # OBJECT or BODY
                    half_width = slide_width / 2
                    # Content placeholder left edge should be roughly at or after the slide midpoint
                    # Allow some tolerance (within 20% of slide width from midpoint)
                    if ph.left >= half_width * 0.8:
                        has_content_right = True
                    else:
                        print(f"  INFO: Content PH left={ph.left}, slide_half={half_width} — not on right half")

            # Award partial credit: 0.10 each for the three sub-checks
            if has_title:
                sub_score += 0.10
                print(f"  PASS: Component 2a — Title placeholder present")
            else:
                print(f"  FAIL: Component 2a — No title placeholder found in 'Feature Slide' layout")

            if has_image_left:
                sub_score += 0.10
                print(f"  PASS: Component 2b — Image/picture placeholder on left half")
            else:
                print(f"  FAIL: Component 2b — No image placeholder on left half found")

            if has_content_right:
                sub_score += 0.10
                print(f"  PASS: Component 2c — Content/text placeholder on right half")
            else:
                print(f"  FAIL: Component 2c — No content placeholder on right half found")

            if sub_score > 0:
                print(f"PASS: Component 2 — Layout placeholders ({sub_score:.2f} pts)")
                total_score += sub_score
            else:
                print(f"FAIL: Component 2 — No valid placeholders found in 'Feature Slide' layout")

        except Exception as e:
            print(f"ERROR: Component 2 — {e}")
    else:
        print("SKIP: Component 2 — 'Feature Slide' layout does not exist")

    # -------------------------------------------------------------------
    # Component 3: Slides 3, 5, 7 use 'Feature Slide' layout (0.40)
    #   - Also verify other slides are NOT changed to 'Feature Slide'
    # -------------------------------------------------------------------
    try:
        # Expected layouts for each slide (1-indexed)
        # Initial state: slides 3,5,7 are 'Blank', others are various layouts
        target_slides = {3, 5, 7}  # 1-indexed slides that should use 'Feature Slide'

        correct_targets = 0
        total_targets = len(target_slides)
        other_slides_ok = True

        for si in range(len(prs.slides)):
            slide = prs.slides[si]
            slide_num = si + 1
            layout_name = slide.slide_layout.name

            if slide_num in target_slides:
                if layout_name == 'Feature Slide':
                    correct_targets += 1
                    print(f"  PASS: Slide {slide_num} uses 'Feature Slide' layout")
                else:
                    print(f"  FAIL: Slide {slide_num} uses '{layout_name}' instead of 'Feature Slide'")
            else:
                # Other slides should NOT use 'Feature Slide'
                if layout_name == 'Feature Slide':
                    other_slides_ok = False
                    print(f"  FAIL: Slide {slide_num} incorrectly uses 'Feature Slide' (should be unchanged)")

        # Score: proportional to how many target slides are correct
        # 0.30 for target slides, 0.10 for other slides being unaffected
        # The "other slides unaffected" bonus only applies if at least one target slide is correct
        # (otherwise it rewards a pre-existing condition)
        target_score = (correct_targets / total_targets) * 0.30 if total_targets > 0 else 0.0
        other_score = 0.10 if (other_slides_ok and correct_targets > 0) else 0.0
        comp3_score = target_score + other_score

        if comp3_score > 0:
            print(f"PASS: Component 3 — Slide layout assignments ({comp3_score:.2f} pts)")
            total_score += comp3_score
        else:
            print(f"FAIL: Component 3 — No target slides use 'Feature Slide' and/or other slides affected")

    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
