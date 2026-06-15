"""
Reward Script: Set bulleted list on slide 4 to 24pt with 1.5 line spacing
Task ID: impress_teach_012
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): All 5 bullet runs have font size 24pt (304800 EMU)
  Component 2 (0.5): All 5 bullet paragraphs have 150% line spacing
"""

import os
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_012'
EXPECTED_FONT_SIZE = 304800  # 24pt in EMU
EXPECTED_LINE_SPACING_PCT = 150000  # 150% = 1.5x, stored as 150000 in spcPct val
EXPECTED_BULLET_COUNT = 5


def get_line_spacing_pct(para):
    """Extract line spacing percentage from paragraph XML. Returns None if not set."""
    pPr = para._p.find(qn('a:pPr'))
    if pPr is None:
        return None
    lnSpc = pPr.find(qn('a:lnSpc'))
    if lnSpc is None:
        return None
    spcPct = lnSpc.find(qn('a:spcPct'))
    if spcPct is not None:
        return int(spcPct.get('val'))
    return None


def persist_app_state(domain):
    """Try to save any unsaved changes in LibreOffice."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # slide 4 (0-indexed)

    # Find the content placeholder (the bulleted list shape, not the title)
    content_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name != "Title 1":
            # Look for the shape with the bullet items
            non_empty_paras = [p for p in shape.text_frame.paragraphs if p.text.strip()]
            if len(non_empty_paras) >= 3:  # at least several bullet items
                content_shape = shape
                break

    if content_shape is None:
        print("FAIL: No content shape with bullet list found on slide 4")
        print("REWARD: 0.0")
        return 0.0

    # Collect non-empty paragraphs (the bullet items)
    bullet_paras = [p for p in content_shape.text_frame.paragraphs if p.text.strip()]
    print(f"INFO: Found {len(bullet_paras)} bullet paragraphs on slide 4")

    # Component 1: Font size is 24pt for all bullet runs (0.5 points)
    try:
        font_pass_count = 0
        font_total_runs = 0
        for para in bullet_paras:
            runs = [r for r in para.runs if (r.text or "").strip()]
            for run in runs:
                font_total_runs += 1
                if run.font.size == EXPECTED_FONT_SIZE:
                    font_pass_count += 1
                else:
                    actual_pt = run.font.size / 12700 if run.font.size else None
                    print(f"FAIL: Run '{run.text[:30]}' has size {actual_pt}pt, expected 24pt")

        if font_total_runs > 0 and font_pass_count == font_total_runs:
            print(f"PASS: Component 1 -- All {font_total_runs} runs are 24pt (0.5 pts)")
            total_score += 0.5
        elif font_total_runs > 0:
            # Partial: proportional credit
            partial = 0.5 * (font_pass_count / font_total_runs)
            print(f"PARTIAL: Component 1 -- {font_pass_count}/{font_total_runs} runs are 24pt ({partial:.2f} pts)")
            total_score += partial
        else:
            print("FAIL: Component 1 -- No runs found in bullet paragraphs")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Line spacing is 150% for all bullet paragraphs (0.5 points)
    try:
        spacing_pass_count = 0
        for para in bullet_paras:
            actual_spacing = get_line_spacing_pct(para)
            if actual_spacing is not None and abs(actual_spacing - EXPECTED_LINE_SPACING_PCT) < 1000:
                # Allow small tolerance (within 1% of 150%)
                spacing_pass_count += 1
            else:
                actual_display = f"{actual_spacing / 1000}%" if actual_spacing else "not set (default)"
                print(f"FAIL: Para '{para.text[:30]}...' has line spacing {actual_display}, expected 150%")

        if len(bullet_paras) > 0 and spacing_pass_count == len(bullet_paras):
            print(f"PASS: Component 2 -- All {len(bullet_paras)} paragraphs have 150% line spacing (0.5 pts)")
            total_score += 0.5
        elif len(bullet_paras) > 0:
            partial = 0.5 * (spacing_pass_count / len(bullet_paras))
            print(f"PARTIAL: Component 2 -- {spacing_pass_count}/{len(bullet_paras)} paragraphs have 150% spacing ({partial:.2f} pts)")
            total_score += partial
        else:
            print("FAIL: Component 2 -- No bullet paragraphs found")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
