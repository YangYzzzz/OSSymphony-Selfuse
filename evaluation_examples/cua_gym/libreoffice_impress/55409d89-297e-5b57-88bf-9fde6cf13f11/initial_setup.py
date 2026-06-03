"""
Initial Setup: Anatomy Lecture presentation with 7 slides.
Slide 4 has a bulleted list with 5 items in 14pt, single line spacing.
Task ID: impress_teach_012
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_012'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_line_spacing(paragraph, spacing_pt):
    """Set line spacing via XML (python-pptx has limited direct support)."""
    pPr = paragraph._p.get_or_add_pPr()
    lnSpc = pPr.find(qn('a:lnSpc'))
    if lnSpc is not None:
        pPr.remove(lnSpc)
    lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
    spcPts = lnSpc.makeelement(qn('a:spcPts'), {'val': str(int(spacing_pt * 100))})
    lnSpc.append(spcPts)
    pPr.append(lnSpc)


def add_text_slide(prs, layout_idx, title_text, body_items, font_size=Pt(18), bold_title=True):
    """Helper to add a slide with title and body text items."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, item in enumerate(body_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.level = 0
            for run in p.runs:
                run.font.size = font_size
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Human Anatomy: Musculoskeletal System"
    slide1.placeholders[1].text = "Dr. Elena Vasquez\nDepartment of Biomedical Sciences\nSpring 2026"

    # --- Slide 2: Course Overview ---
    slide2 = add_text_slide(prs, 1, "Course Overview", [
        "Comprehensive study of skeletal and muscular anatomy",
        "Focus on clinical applications and functional relationships",
        "Weekly lab dissections complementing lecture material",
        "Assessment: midterm (30%), final (40%), lab practical (30%)",
    ])

    # --- Slide 3: Skeletal System Introduction ---
    slide3 = add_text_slide(prs, 1, "The Skeletal System", [
        "206 bones in the adult human skeleton",
        "Axial skeleton: skull, vertebral column, thoracic cage",
        "Appendicular skeleton: limbs and girdles",
        "Functions: support, protection, movement, mineral storage",
        "Bone remodeling is a continuous, lifelong process",
    ])

    # --- Slide 4: Key Learning Objectives (THE BULLETED LIST) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Key Learning Objectives"
    tf4 = slide4.placeholders[1].text_frame
    tf4.clear()

    bullet_items = [
        "Identify the major bones of the axial and appendicular skeleton",
        "Describe the histological structure of compact and spongy bone",
        "Explain the process of endochondral ossification in long bones",
        "Compare and contrast the three types of muscle tissue",
        "Analyze the biomechanics of synovial joint movement",
    ]

    for i, item in enumerate(bullet_items):
        if i == 0:
            p = tf4.paragraphs[0]
        else:
            p = tf4.add_paragraph()
        p.text = item
        p.level = 0
        # Set font to 14pt
        for run in p.runs:
            run.font.size = Pt(14)
        # Single line spacing (use spcPts for explicit control)
        # 14pt * 100 = 1400 (single spacing matches font size in spcPts)
        # We do NOT set lnSpc at all for single spacing (default behavior)

    # --- Slide 5: Muscle Classification ---
    slide5 = add_text_slide(prs, 1, "Muscle Tissue Classification", [
        "Skeletal muscle: voluntary, striated, attached to bones",
        "Cardiac muscle: involuntary, striated, found only in the heart",
        "Smooth muscle: involuntary, non-striated, lines hollow organs",
        "Over 600 skeletal muscles in the human body",
    ])

    # --- Slide 6: Clinical Applications ---
    slide6 = add_text_slide(prs, 1, "Clinical Applications", [
        "Osteoporosis: reduced bone density, increased fracture risk",
        "Muscular dystrophy: progressive degeneration of muscle fibers",
        "Arthritis: inflammation of joints affecting mobility",
        "Imaging techniques: X-ray, MRI, CT for diagnostic evaluation",
        "Rehabilitation protocols for common musculoskeletal injuries",
    ])

    # --- Slide 7: Next Lecture & Readings ---
    slide7 = add_text_slide(prs, 1, "Next Lecture & Required Readings", [
        "Topic: The Nervous System - Central and Peripheral Divisions",
        "Read: Netter's Atlas of Human Anatomy, Chapters 8-10",
        "Review: Lab manual exercises on spinal cord cross-sections",
        "Optional: Virtual dissection module available on course portal",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
