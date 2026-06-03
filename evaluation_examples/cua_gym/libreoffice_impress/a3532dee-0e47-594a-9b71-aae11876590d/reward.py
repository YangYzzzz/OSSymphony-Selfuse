"""
Reward Script: CI/CD Pipeline Diagram on Slide 7
Task ID: impress_ps_033
Domain: libreoffice_impress
Scoring:
  Component 1: Five stage rectangles with correct names (0.30)
  Component 2: Correct fill colors on stage rectangles (0.25)
  Component 3: Five green status indicator circles (0.20)
  Component 4: Four arrow connectors between stages (0.15)
  Component 5: Text formatting - bold, contrasting color, ~14pt (0.10)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt, Emu

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_033'

# Expected stage names in order
STAGE_NAMES = ['Code Commit', 'Build', 'Test', 'Stage Deploy', 'Prod Deploy']

# Expected fill colors per stage (from task instruction)
# blue, yellow, orange, purple, green
STAGE_COLORS = {
    'Code Commit': None,     # blue - we'll check it's a blue-ish color
    'Build': None,           # yellow
    'Test': None,            # orange
    'Stage Deploy': None,    # purple
    'Prod Deploy': None,     # green
}

# Color families: we check that each stage has a color in the correct family
def is_blue(rgb_str):
    """Check if color is blue-ish"""
    r, g, b = int(rgb_str[0:2], 16), int(rgb_str[2:4], 16), int(rgb_str[4:6], 16)
    return b > 150 and b > r and b > g

def is_yellow(rgb_str):
    """Check if color is yellow-ish"""
    r, g, b = int(rgb_str[0:2], 16), int(rgb_str[2:4], 16), int(rgb_str[4:6], 16)
    return r > 180 and g > 150 and b < 100

def is_orange(rgb_str):
    """Check if color is orange-ish"""
    r, g, b = int(rgb_str[0:2], 16), int(rgb_str[2:4], 16), int(rgb_str[4:6], 16)
    return r > 200 and g > 100 and g < 200 and b < 80

def is_purple(rgb_str):
    """Check if color is purple-ish"""
    r, g, b = int(rgb_str[0:2], 16), int(rgb_str[2:4], 16), int(rgb_str[4:6], 16)
    return r > 100 and b > 100 and b > g

def is_green(rgb_str):
    """Check if color is green-ish"""
    r, g, b = int(rgb_str[0:2], 16), int(rgb_str[2:4], 16), int(rgb_str[4:6], 16)
    return g > 130 and g > r and g > b


COLOR_CHECKERS = {
    'Code Commit': is_blue,
    'Build': is_yellow,
    'Test': is_orange,
    'Stage Deploy': is_purple,
    'Prod Deploy': is_green,
}


def get_shape_fill_rgb(shape):
    """Get fill color RGB string from a shape, or None."""
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb)
    except:
        pass
    return None


def get_shape_text(shape):
    """Get all text from a shape."""
    if hasattr(shape, 'text'):
        return shape.text.strip()
    return ""


def verify_task(file_path):
    """
    Verify CI/CD pipeline diagram on slide 7.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Check we have at least 7 slides
    if len(prs.slides) < 7:
        print(f"FAIL: Need at least 7 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[6]  # Slide 7 (0-indexed)
    shapes = list(slide.shapes)

    # Classify shapes on slide 7 (excluding original title/textbox)
    # We need to find: rounded rectangles with stage text, ovals (indicators), arrows
    stage_shapes = []  # shapes containing stage name text
    oval_shapes = []   # small circles (status indicators)
    arrow_shapes = []  # arrow connectors

    for s in shapes:
        name_lower = s.name.lower() if s.name else ""
        text = get_shape_text(s)

        # Skip the original title placeholder and CI/CD Pipeline textbox
        if s.shape_type == 14:  # PLACEHOLDER
            continue
        if s.shape_type == 17:  # TEXT_BOX - the original title
            continue

        # Identify by content
        if text and any(stage.lower() == text.lower() for stage in STAGE_NAMES):
            stage_shapes.append(s)
        elif 'arrow' in name_lower or 'arrow' in str(getattr(s, 'auto_shape_type', '')).lower():
            arrow_shapes.append(s)
        elif 'oval' in name_lower or 'circle' in name_lower:
            oval_shapes.append(s)
        else:
            # Check if it's a small shape (likely indicator) vs larger shape
            if s.width and s.height:
                # Small shapes (< 1 inch in both dims) are likely indicators
                if s.width < 914400 and s.height < 914400 and not text:
                    oval_shapes.append(s)
                elif text:
                    # Could be a stage shape with slightly different text
                    stage_shapes.append(s)

    print(f"Found: {len(stage_shapes)} stage shapes, {len(oval_shapes)} indicators, {len(arrow_shapes)} arrows")

    # Component 1: Five stage rectangles with correct names (0.30 points)
    try:
        found_stages = set()
        for s in stage_shapes:
            text = get_shape_text(s).strip()
            for stage_name in STAGE_NAMES:
                if text.lower() == stage_name.lower():
                    found_stages.add(stage_name)
                    break

        stage_count = len(found_stages)
        if stage_count == 5:
            print(f"PASS: Component 1 - All 5 stage rectangles found: {found_stages} (0.30 pts)")
            total_score += 0.30
        elif stage_count > 0:
            partial = round(0.30 * stage_count / 5, 2)
            print(f"PARTIAL: Component 1 - {stage_count}/5 stages found: {found_stages} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 - No stage rectangles found with expected names")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Correct fill colors on stage rectangles (0.25 points)
    try:
        correct_colors = 0
        for s in stage_shapes:
            text = get_shape_text(s).strip()
            rgb = get_shape_fill_rgb(s)
            if rgb is None:
                continue
            for stage_name, checker in COLOR_CHECKERS.items():
                if text.lower() == stage_name.lower():
                    if checker(rgb):
                        correct_colors += 1
                        print(f"  Color OK: {stage_name} = {rgb}")
                    else:
                        print(f"  Color MISMATCH: {stage_name} = {rgb}")
                    break

        if correct_colors == 5:
            print(f"PASS: Component 2 - All 5 stages have correct colors (0.25 pts)")
            total_score += 0.25
        elif correct_colors > 0:
            partial = round(0.25 * correct_colors / 5, 2)
            print(f"PARTIAL: Component 2 - {correct_colors}/5 correct colors ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 - No stages have correct fill colors")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Five green status indicator circles (0.20 points)
    try:
        green_indicators = 0
        for s in oval_shapes:
            rgb = get_shape_fill_rgb(s)
            if rgb and is_green(rgb):
                green_indicators += 1

        if green_indicators >= 5:
            print(f"PASS: Component 3 - {green_indicators} green indicator circles found (0.20 pts)")
            total_score += 0.20
        elif green_indicators > 0:
            partial = round(0.20 * min(green_indicators, 5) / 5, 2)
            print(f"PARTIAL: Component 3 - {green_indicators}/5 green indicators ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 - No green indicator circles found")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Four arrow connectors between stages (0.15 points)
    try:
        arrow_count = len(arrow_shapes)
        if arrow_count >= 4:
            print(f"PASS: Component 4 - {arrow_count} arrow connectors found (0.15 pts)")
            total_score += 0.15
        elif arrow_count > 0:
            partial = round(0.15 * arrow_count / 4, 2)
            print(f"PARTIAL: Component 4 - {arrow_count}/4 arrows ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 - No arrow connectors found")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Text formatting - bold, contrasting color, ~14pt (0.10 points)
    try:
        well_formatted = 0
        for s in stage_shapes:
            text = get_shape_text(s)
            if not text:
                continue
            for stage_name in STAGE_NAMES:
                if text.lower() != stage_name.lower():
                    continue
                for p in s.text_frame.paragraphs:
                    for r in p.runs:
                        if not r.text.strip():
                            continue
                        is_bold = r.font.bold is True or r.font.bold is None  # None can inherit
                        # Check font size is approximately 14pt (177800 EMU = 14pt)
                        size_ok = False
                        if r.font.size is not None:
                            pt_size = r.font.size / 12700  # EMU to pt
                            size_ok = 10 <= pt_size <= 20  # reasonable range
                        else:
                            size_ok = (r.font.size is None)  # inherited size is acceptable
                        # Check contrasting text color (white or dark on colored bg)
                        try:
                            color_rgb = str(r.font.color.rgb) if r.font.color.type is not None else None
                        except:
                            color_rgb = None
                        has_contrast = color_rgb is not None  # any explicit color = intentional contrast

                        if is_bold and size_ok and has_contrast:
                            well_formatted += 1
                        break  # only check first non-empty run
                break

        if well_formatted >= 5:
            print(f"PASS: Component 5 - All 5 stages have proper text formatting (0.10 pts)")
            total_score += 0.10
        elif well_formatted > 0:
            partial = round(0.10 * well_formatted / 5, 2)
            print(f"PARTIAL: Component 5 - {well_formatted}/5 stages formatted ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 - Stage text lacks proper formatting")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
