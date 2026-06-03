"""
Initial Setup: Create Tech_Architecture.pptx with 8 slides, slide 7 has title 'CI/CD Pipeline' but empty content.
Task ID: impress_ps_033
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_033'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_body(slide, title_text, body_lines):
    """Helper to set title and body text on a slide with Title+Content layout."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    # Find body placeholder (idx 1)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.text = body_lines[0] if body_lines else ""
            for line in body_lines[1:]:
                p = tf.add_paragraph()
                p.text = line
                p.level = 0
            break


def create_initial():
    prs = Presentation()
    # Standard widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Tech Architecture Overview"
    slide1.placeholders[1].text = "Enterprise Platform — Q1 2026 Review"

    # --- Slide 2: System Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide2, "System Overview", [
        "Microservices architecture with 47 active services",
        "Kubernetes orchestration across 3 cloud regions",
        "Event-driven communication via Apache Kafka",
        "Average latency: 23ms at p99",
        "99.97% uptime over the last 12 months",
    ])

    # --- Slide 3: Frontend Architecture ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide3, "Frontend Architecture", [
        "React 18 with TypeScript for web applications",
        "Next.js server-side rendering for SEO-critical pages",
        "Design system: 120+ reusable components",
        "CDN: Cloudflare with edge caching",
        "Mobile: React Native shared codebase (iOS/Android)",
    ])

    # --- Slide 4: Backend Services ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide4, "Backend Services", [
        "Go and Python services behind API Gateway",
        "PostgreSQL primary database with read replicas",
        "Redis caching layer — 94% cache hit rate",
        "Elasticsearch for full-text search (2.3B documents)",
        "gRPC for inter-service communication",
    ])

    # --- Slide 5: Data Platform ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide5, "Data Platform", [
        "Snowflake data warehouse — 18TB structured data",
        "Apache Airflow for ETL orchestration (340 DAGs)",
        "Real-time streaming: Kafka Streams + Flink",
        "ML model serving via SageMaker endpoints",
        "Data governance with Apache Atlas metadata catalog",
    ])

    # --- Slide 6: Security & Compliance ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide6, "Security & Compliance", [
        "Zero-trust network architecture",
        "OAuth 2.0 + OIDC authentication layer",
        "SOC 2 Type II and ISO 27001 certified",
        "Vault for secrets management (HashiCorp Vault)",
        "Automated vulnerability scanning in CI pipeline",
    ])

    # --- Slide 7: CI/CD Pipeline (TASK TARGET — title only, no content shapes) ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title as a textbox
    txBox = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "CI/CD Pipeline"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 8: Roadmap & Next Steps ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide8, "Roadmap & Next Steps", [
        "Q2: Service mesh migration (Istio rollout)",
        "Q3: Multi-region active-active deployment",
        "Q4: Edge computing pilot for IoT workloads",
        "Ongoing: Platform team hiring (3 SREs, 2 architects)",
        "Target: Sub-10ms p99 latency by end of year",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
