"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 114 looks a bit plain. In LibreOffice Impress, how can I give just the very first word of the title a dotted underline and leave the rest of the title untouched?
Generated: 2025-09-11 00:18:15
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
from pptx import Presentation
from pptx.enum.text import MSO_TEXT_UNDERLINE_TYPE as MSO_UNDERLINE


def verify_dotted_first_word_title(file_path: str) -> float:
    """Verify that on slide 114 of the given PPTX file the very first word
    of the title has a dotted underline while the rest of the title words
    remain without any underline styling.

    Scoring (progressive):
        0.2 – Slide 114 exists
        0.1 – Title shape located
        0.3 – First word has dotted underline (light or heavy)
        0.4 – Remaining words have NO underline
        -----
        1.0 – Perfect completion
    """

    print(f"Verifying presentation: {file_path}")
    score = 0.0
    max_score = 1.0

    # ------------------------------------------------------------------
    # 1. File existence & loading  (NO POINTS – prerequisite only)
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File not found")
        return 0.0
    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded presentation with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Could not load PPTX: {e}")
        return 0.0

    # ------------------------------------------------------------------
    # 2. Ensure slide 114 exists  (0.2 pts)
    # ------------------------------------------------------------------
    if len(prs.slides) >= 114:
        print("✓ Slide 114 exists")
        score += 0.2
        slide = prs.slides[113]  # zero-based index
    else:
        print("✗ Slide 114 does not exist")
        return score

    # ------------------------------------------------------------------
    # 3. Locate the title shape  (0.1 pts)
    # ------------------------------------------------------------------
    title_shape = None
    for shape in slide.shapes:
        # Placeholder type 1 == TITLE
        if shape.is_placeholder and shape.placeholder_format.type == 1:
            title_shape = shape
            break

    # Fallbacks if placeholder detection failed
    if title_shape is None:
        for shape in slide.shapes:
            if shape.name.lower().startswith("title") and shape.has_text_frame:
                title_shape = shape
                break
    if title_shape is None and any(s.has_text_frame for s in slide.shapes):
        title_shape = next(s for s in slide.shapes if s.has_text_frame)

    if title_shape is None:
        print("✗ No title text found on slide 114")
        return score
    else:
        print(f"✓ Found title shape: {title_shape.name}")
        score += 0.1

    # ------------------------------------------------------------------
    # 4. Analyse underline styling of the first and subsequent words
    # ------------------------------------------------------------------
    tf = title_shape.text_frame
    if not tf.text.strip():
        print("✗ Title text is empty")
        return score

    paragraph = tf.paragraphs[0]
    print("Title paragraph text:", paragraph.text)

    if not paragraph.runs:
        print("✗ Title paragraph has no runs")
        return score

    dotted_values = {MSO_UNDERLINE.DOTTED_LINE, MSO_UNDERLINE.DOTTED_HEAVY_LINE}

    # ---- First word/run check (0.3 pts) ----
    first_run = paragraph.runs[0]
    print("First run text:", first_run.text, " underline:", first_run.font.underline)

    if first_run.font.underline in dotted_values:
        print("✓ First word has dotted underline")
        score += 0.3
    else:
        print("✗ First word does not have dotted underline")

    # ---- Remaining runs check (0.4 pts) ----
    other_runs = paragraph.runs[1:]
    other_ok = True
    for run in other_runs:
        if run.font.underline in dotted_values or (
            run.font.underline is not None and run.font.underline != MSO_UNDERLINE.NONE
        ):
            other_ok = False
            print("✗ Other run underlined:", run.text, run.font.underline)
            break

    if other_runs:
        if other_ok:
            print("✓ Other words are not underlined")
            score += 0.4
    else:
        # Title contains a single word -> nothing else to verify
        print("ℹ️ Title only has one run (treated as no extra underlines)")
        score += 0.4

    # ------------------------------------------------------------------
    # 5. Final score & output
    # ------------------------------------------------------------------
    final_score = min(score, max_score)
    print(f"Total score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score


if __name__ == "__main__":
    FILE_PATH = (
        "/home/user/slide_114_looks_a_bit_plain_in_libreoffice_impress_how_can_i_give_just_the_very_first_word_of_the_ti_golden.pptx"
    )
    verify_dotted_first_word_title(FILE_PATH)

