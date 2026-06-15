"""
Initial Setup: Insert a bulleted list text box on slide 3
Task ID: impress_tm_070
Domain: libreoffice_impress

Creates a 5-slide presentation. Slide 3 ('Project Phases') has only a title,
no bulleted list text box — the agent must add it.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_070'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Project Plan 2025"
    slide1.placeholders[1].text = "Prepared by the Strategy Team"

    # --- Slide 2: Project Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Project Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "This project aims to modernize our internal tooling infrastructure."
    p2 = body2.add_paragraph()
    p2.text = "The initiative spans Q1-Q3 2025 with a total budget of $1.2M."
    p3 = body2.add_paragraph()
    p3.text = "Key stakeholders include Engineering, Product, and Operations teams."

    # --- Slide 3: Project Phases (TITLE ONLY - no bulleted list) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add only the title as a text box at the top
    title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Project Phases"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 4: Timeline ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Timeline"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Q1 2025: Planning and requirements gathering"
    for item in [
        "Q2 2025: Core development and integration",
        "Q3 2025: User acceptance testing and deployment",
        "Q4 2025: Post-launch monitoring and optimization",
    ]:
        p = body4.add_paragraph()
        p.text = item

    # --- Slide 5: Budget Summary ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Budget Summary"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Personnel costs: $680,000"
    for item in [
        "Software licenses: $185,000",
        "Infrastructure: $220,000",
        "Contingency fund: $115,000",
        "Total: $1,200,000",
    ]:
        p = body5.add_paragraph()
        p.text = item

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
