"""
Reward Script: Team slide with 2x2 grid of member cards on slide 7
Task ID: impress_ps_005
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Four circular shapes (photo placeholders) exist on slide 7
  Component 2 (0.30): Four name text boxes with correct names, bold, ~16pt
  Component 3 (0.25): Four title text boxes with correct titles, ~12pt, not bold
  Component 4 (0.15): 2x2 grid layout (two distinct rows and two distinct columns)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_005'

# Expected team members: name -> title
TEAM_MEMBERS = {
    'Sarah Chen': 'CEO',
    'Marcus Lee': 'CTO',
    'Priya Patel': 'VP Engineering',
    'James Wilson': 'CFO',
}


def persist_app_state(domain):
    """Try to save any unsaved LibreOffice state via Ctrl+S."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print(f"PERSIST: ctrl+s sent for {domain}")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def get_all_shapes_recursive(slide):
    """Get all shapes including those nested in groups."""
    results = []
    for shape in slide.shapes:
        results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.append(sub)
                if hasattr(sub, 'shapes'):
                    for ssub in sub.shapes:
                        results.append(ssub)
    return results


def normalize_text(t):
    """Normalize text for comparison: strip and collapse whitespace."""
    return ' '.join((t or '').split()).strip()


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 7 slides
    if len(prs.slides) < 7:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 7")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[6]  # Slide 7 (0-indexed)
    all_shapes = get_all_shapes_recursive(slide)

    # Identify shape types on the slide (excluding the existing "Our Team" title)
    ovals = []
    text_shapes = []
    for shape in all_shapes:
        # Oval/circle shapes (AUTO_SHAPE type 1, or freeform circles)
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # Check if it's roughly circular (width ~= height)
            if shape.width > 0 and shape.height > 0:
                ratio = shape.width / shape.height
                if 0.8 <= ratio <= 1.25:  # roughly circular
                    ovals.append(shape)
        # Also check for FREEFORM or OVAL types
        elif shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            if shape.width > 0 and shape.height > 0:
                ratio = shape.width / shape.height
                if 0.8 <= ratio <= 1.25:
                    ovals.append(shape)

        # Text boxes (excluding the title "Our Team")
        if hasattr(shape, 'text_frame') and hasattr(shape, 'text'):
            txt = normalize_text(shape.text)
            if txt and txt.lower() != 'our team':
                text_shapes.append(shape)

    # Component 1: Four circular shapes (photo placeholders) (0.30 points)
    # Task requires circular placeholders for photos, ~3cm diameter
    try:
        circle_count = len(ovals)
        if circle_count >= 4:
            # Verify they are roughly 3cm diameter (1080000 EMU = 3cm)
            # Allow tolerance: 1.5cm to 5cm diameter
            valid_circles = [o for o in ovals if 540000 <= o.width <= 1800000]
            if len(valid_circles) >= 4:
                print(f"PASS: Component 1 - Found {len(valid_circles)} circular photo placeholders (0.30 pts)")
                total_score += 0.30
            else:
                print(f"PARTIAL: Component 1 - Found {circle_count} circles but only {len(valid_circles)} in valid size range")
                total_score += 0.15
        elif circle_count > 0:
            print(f"PARTIAL: Component 1 - Found only {circle_count}/4 circular shapes")
            total_score += 0.10 * min(circle_count, 4) / 4
        else:
            print(f"FAIL: Component 1 - No circular photo placeholder shapes found on slide 7")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Four name text boxes with correct names, bold, ~16pt (0.30 points)
    try:
        found_names = {}
        for shape in text_shapes:
            txt = normalize_text(shape.text)
            for name in TEAM_MEMBERS:
                if name.lower() in txt.lower():
                    # Check formatting
                    is_bold = False
                    size_ok = False
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run_txt = normalize_text(run.text)
                            if run_txt and name.lower() in run_txt.lower():
                                if run.font.bold is True:
                                    is_bold = True
                                if run.font.size is not None:
                                    pt_size = run.font.size / 12700
                                    if 13 <= pt_size <= 20:  # ~16pt with tolerance
                                        size_ok = True
                    found_names[name] = {'bold': is_bold, 'size_ok': size_ok}

        names_found = len(found_names)
        names_with_bold = sum(1 for v in found_names.values() if v['bold'])
        names_with_size = sum(1 for v in found_names.values() if v['size_ok'])

        if names_found >= 4 and names_with_bold >= 4:
            print(f"PASS: Component 2 - All 4 names found with bold formatting (0.30 pts)")
            total_score += 0.30
        elif names_found >= 4:
            # Names present but not all bold
            partial = 0.15 + 0.15 * (names_with_bold / 4)
            print(f"PARTIAL: Component 2 - All 4 names found, {names_with_bold}/4 bold ({partial:.2f} pts)")
            total_score += partial
        elif names_found > 0:
            partial = 0.30 * (names_found / 4) * 0.5
            print(f"PARTIAL: Component 2 - Found {names_found}/4 names ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 - No team member names found on slide 7")
            print(f"  Available text shapes: {[normalize_text(s.text) for s in text_shapes]}")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Four title text boxes with correct titles, ~12pt, not bold (0.25 points)
    try:
        titles = list(TEAM_MEMBERS.values())  # CEO, CTO, VP Engineering, CFO
        found_titles = {}
        for shape in text_shapes:
            txt = normalize_text(shape.text)
            for title in titles:
                if title.lower() == txt.lower():
                    # Check formatting: should be ~12pt and not bold
                    not_bold = False
                    size_ok = False
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run_txt = normalize_text(run.text)
                            if run_txt and title.lower() in run_txt.lower():
                                # Bold should be False or None (not set = inherit = not bold)
                                if run.font.bold is not True:
                                    not_bold = True
                                if run.font.size is not None:
                                    pt_size = run.font.size / 12700
                                    if 10 <= pt_size <= 14:  # ~12pt with tolerance
                                        size_ok = True
                    found_titles[title] = {'not_bold': not_bold, 'size_ok': size_ok}

        titles_found = len(found_titles)
        if titles_found >= 4:
            print(f"PASS: Component 3 - All 4 titles found (0.25 pts)")
            total_score += 0.25
        elif titles_found > 0:
            partial = 0.25 * (titles_found / 4)
            print(f"PARTIAL: Component 3 - Found {titles_found}/4 titles ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 - No team member titles found on slide 7")
            print(f"  Available text shapes: {[normalize_text(s.text) for s in text_shapes]}")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: 2x2 grid layout (0.15 points)
    # Check that shapes form approximately 2 rows and 2 columns
    try:
        if len(ovals) >= 4:
            # Use oval positions to determine grid layout
            tops = sorted(set(o.top for o in ovals[:4]))
            lefts = sorted(set(o.left for o in ovals[:4]))

            # Cluster tops into rows (within 15% of slide height tolerance)
            slide_height = prs.slide_height
            row_clusters = []
            for t in sorted(o.top for o in ovals[:4]):
                placed = False
                for cluster in row_clusters:
                    if abs(t - cluster[0]) < slide_height * 0.15:
                        cluster.append(t)
                        placed = True
                        break
                if not placed:
                    row_clusters.append([t])

            # Cluster lefts into columns
            slide_width = prs.slide_width
            col_clusters = []
            for l in sorted(o.left for o in ovals[:4]):
                placed = False
                for cluster in col_clusters:
                    if abs(l - cluster[0]) < slide_width * 0.15:
                        cluster.append(l)
                        placed = True
                        break
                if not placed:
                    col_clusters.append([l])

            num_rows = len(row_clusters)
            num_cols = len(col_clusters)

            if num_rows == 2 and num_cols == 2:
                print(f"PASS: Component 4 - 2x2 grid layout detected ({num_rows} rows, {num_cols} cols) (0.15 pts)")
                total_score += 0.15
            elif num_rows >= 2 or num_cols >= 2:
                print(f"PARTIAL: Component 4 - Layout has {num_rows} rows and {num_cols} cols, expected 2x2 (0.08 pts)")
                total_score += 0.08
            else:
                print(f"FAIL: Component 4 - Not a 2x2 grid: {num_rows} rows, {num_cols} cols")
        else:
            print(f"FAIL: Component 4 - Not enough circular shapes ({len(ovals)}) to verify grid layout")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
