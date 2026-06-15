"""
Initial Setup: Build a startup pitch presentation with 9 slides.
Slide 7 has title 'Our Team' but is otherwise empty (no team cards yet).
Task ID: impress_ps_005
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_005'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_paragraph(tf, text, font_size=Pt(18), bold=False, alignment=PP_ALIGN.LEFT,
                  color=None, font_name="Arial"):
    """Helper to set text on first paragraph of a text frame."""
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    for run in p.runs:
        run.font.size = font_size
        run.font.bold = bold
        run.font.name = font_name
        if color:
            run.font.color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=Pt(18),
                 bold=False, alignment=PP_ALIGN.LEFT, color=None, font_name="Arial"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    set_paragraph(tf, text, font_size, bold, alignment, color, font_name)
    return txBox


def create_initial():
    prs = Presentation()
    # Standard 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Dark background
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)
    add_text_box(slide1, Inches(1.5), Inches(1.5), Inches(10), Inches(2),
                 "NovaTech Solutions", Pt(44), True, PP_ALIGN.CENTER,
                 RGBColor(0xFF, 0xFF, 0xFF))
    add_text_box(slide1, Inches(2), Inches(3.8), Inches(9), Inches(1.2),
                 "Revolutionizing Supply Chain Intelligence with AI",
                 Pt(22), False, PP_ALIGN.CENTER, RGBColor(0xA0, 0xC4, 0xE8))
    add_text_box(slide1, Inches(4), Inches(5.5), Inches(5), Inches(0.8),
                 "Series A Pitch | Q2 2025", Pt(16), False, PP_ALIGN.CENTER,
                 RGBColor(0x80, 0x80, 0x80))

    # --- Slide 2: Problem ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide2, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "The Problem", Pt(36), True, PP_ALIGN.LEFT,
                 RGBColor(0x0D, 0x1B, 0x2A))
    problems = [
        "Global supply chains lose $1.8 trillion annually to inefficiencies",
        "78% of enterprises still rely on spreadsheet-based demand forecasting",
        "Average lead time variability has increased 42% since 2020",
        "Real-time visibility exists for only 6% of supply chain nodes"
    ]
    y_pos = Inches(1.8)
    for prob in problems:
        add_text_box(slide2, Inches(1.2), y_pos, Inches(10.5), Inches(0.9),
                     f"  {prob}", Pt(18), False, PP_ALIGN.LEFT,
                     RGBColor(0x33, 0x33, 0x33))
        y_pos += Inches(1.2)

    # --- Slide 3: Solution ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide3, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Our Solution", Pt(36), True, PP_ALIGN.LEFT,
                 RGBColor(0x0D, 0x1B, 0x2A))
    add_text_box(slide3, Inches(1), Inches(1.8), Inches(11), Inches(1.5),
                 "NovaTech's AI-powered platform provides end-to-end supply chain visibility, "
                 "predictive analytics, and autonomous optimization. Our proprietary models "
                 "reduce forecasting errors by 67% and cut logistics costs by 23%.",
                 Pt(18), False, PP_ALIGN.LEFT, RGBColor(0x33, 0x33, 0x33))
    # Feature boxes
    features = ["Predictive Demand Engine", "Real-time Visibility Graph", "Autonomous Routing"]
    for i, feat in enumerate(features):
        left = Inches(1 + i * 3.8)
        add_text_box(slide3, left, Inches(4.2), Inches(3.2), Inches(1.5),
                     feat, Pt(16), True, PP_ALIGN.CENTER,
                     RGBColor(0x1A, 0x73, 0xE8))

    # --- Slide 4: Market Opportunity ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide4, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Market Opportunity", Pt(36), True, PP_ALIGN.LEFT,
                 RGBColor(0x0D, 0x1B, 0x2A))
    market_data = [
        "Total Addressable Market: $54.7B by 2028 (CAGR 11.2%)",
        "Serviceable Addressable Market: $12.3B (mid-market enterprises)",
        "Serviceable Obtainable Market: $1.8B (North America + Europe)",
        "Current penetration of AI in supply chain: < 5%"
    ]
    y_pos = Inches(1.8)
    for item in market_data:
        add_text_box(slide4, Inches(1.2), y_pos, Inches(10.5), Inches(0.9),
                     f"  {item}", Pt(18), False, PP_ALIGN.LEFT,
                     RGBColor(0x33, 0x33, 0x33))
        y_pos += Inches(1.2)

    # --- Slide 5: Traction ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide5, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Traction & Milestones", Pt(36), True, PP_ALIGN.LEFT,
                 RGBColor(0x0D, 0x1B, 0x2A))
    milestones = [
        "$2.4M ARR with 340% YoY growth",
        "47 enterprise customers including 3 Fortune 500 companies",
        "Net Revenue Retention: 138%",
        "Average contract value: $52K (up from $31K in 2024)",
        "SOC 2 Type II certified | GDPR compliant"
    ]
    y_pos = Inches(1.8)
    for ms in milestones:
        add_text_box(slide5, Inches(1.2), y_pos, Inches(10.5), Inches(0.8),
                     f"  {ms}", Pt(17), False, PP_ALIGN.LEFT,
                     RGBColor(0x33, 0x33, 0x33))
        y_pos += Inches(1.0)

    # --- Slide 6: Business Model ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide6, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Business Model", Pt(36), True, PP_ALIGN.LEFT,
                 RGBColor(0x0D, 0x1B, 0x2A))
    model_items = [
        "SaaS subscription with usage-based pricing tiers",
        "Starter: $999/mo (up to 50 supply nodes)",
        "Professional: $4,999/mo (up to 500 nodes + advanced analytics)",
        "Enterprise: Custom pricing (unlimited nodes + dedicated support)",
        "Gross margin: 82% | CAC payback: 9 months"
    ]
    y_pos = Inches(1.8)
    for item in model_items:
        add_text_box(slide6, Inches(1.2), y_pos, Inches(10.5), Inches(0.8),
                     f"  {item}", Pt(17), False, PP_ALIGN.LEFT,
                     RGBColor(0x33, 0x33, 0x33))
        y_pos += Inches(1.0)

    # --- Slide 7: Our Team (EMPTY content area - just title) ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide7, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Our Team", Pt(36), True, PP_ALIGN.LEFT,
                 RGBColor(0x0D, 0x1B, 0x2A))
    # Content area is intentionally empty - the task is to add team cards here

    # --- Slide 8: Competitive Advantage ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide8, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Competitive Advantage", Pt(36), True, PP_ALIGN.LEFT,
                 RGBColor(0x0D, 0x1B, 0x2A))
    advantages = [
        "Proprietary transformer model trained on 14B supply chain events",
        "3x faster inference than competing solutions",
        "Patent-pending anomaly detection algorithm",
        "Integrates with 200+ ERP, WMS, and TMS platforms out of the box"
    ]
    y_pos = Inches(1.8)
    for adv in advantages:
        add_text_box(slide8, Inches(1.2), y_pos, Inches(10.5), Inches(0.9),
                     f"  {adv}", Pt(18), False, PP_ALIGN.LEFT,
                     RGBColor(0x33, 0x33, 0x33))
        y_pos += Inches(1.2)

    # --- Slide 9: The Ask ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide9.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)
    add_text_box(slide9, Inches(2), Inches(1.5), Inches(9), Inches(1.5),
                 "The Ask", Pt(40), True, PP_ALIGN.CENTER,
                 RGBColor(0xFF, 0xFF, 0xFF))
    add_text_box(slide9, Inches(2), Inches(3.2), Inches(9), Inches(1.5),
                 "Raising $15M Series A to scale go-to-market, expand engineering, "
                 "and enter APAC markets by Q4 2025.",
                 Pt(20), False, PP_ALIGN.CENTER, RGBColor(0xA0, 0xC4, 0xE8))
    add_text_box(slide9, Inches(3), Inches(5.5), Inches(7), Inches(0.8),
                 "contact@novatechsolutions.ai  |  novatechsolutions.ai",
                 Pt(16), False, PP_ALIGN.CENTER, RGBColor(0x80, 0x80, 0x80))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open the file in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
