"""
Initial Setup: Create a 4-slide presentation with a flat bulleted list on slide 2
Task ID: impstruct_030
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impstruct_030'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Project Roadmap Q3 2025"
    slide1.placeholders[1].text = "Prepared by the Strategy & Operations Team"

    # --- Slide 2: Deliverables (flat bulleted list, 6 items, default bullets) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Deliverables"

    # Use the content placeholder (index 1)
    body = slide2.placeholders[1]
    tf = body.text_frame
    tf.clear()

    items = [
        "Platform Migration",
        "Database schema redesign for PostgreSQL 16",
        "API endpoint compatibility layer deployment",
        "Security Hardening",
        "Multi-factor authentication rollout across all services",
        "Penetration testing report and remediation plan",
    ]

    for i, item_text in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item_text
        p.level = 0  # All flat, no hierarchy
        p.font.size = Pt(18)

    # --- Slide 3: Timeline ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Timeline"

    body3 = slide3.placeholders[1]
    tf3 = body3.text_frame
    tf3.clear()

    timeline_items = [
        "Phase 1 (Jul 1 - Jul 31): Requirements gathering and architecture review",
        "Phase 2 (Aug 1 - Aug 31): Core implementation and unit testing",
        "Phase 3 (Sep 1 - Sep 22): Integration testing and staging deployment",
        "Phase 4 (Sep 23 - Sep 30): Production rollout and monitoring",
    ]

    for i, item_text in enumerate(timeline_items):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = item_text
        p.level = 0
        p.font.size = Pt(16)

    # --- Slide 4: Budget Overview ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Budget Overview"

    body4 = slide4.placeholders[1]
    tf4 = body4.text_frame
    tf4.clear()

    budget_items = [
        "Total Allocated Budget: $485,000",
        "Infrastructure Costs: $210,000 (cloud hosting, licenses)",
        "Personnel: $195,000 (contractors and overtime)",
        "Contingency Reserve: $80,000",
    ]

    for i, item_text in enumerate(budget_items):
        if i == 0:
            p = tf4.paragraphs[0]
        else:
            p = tf4.add_paragraph()
        p.text = item_text
        p.level = 0
        p.font.size = Pt(16)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
