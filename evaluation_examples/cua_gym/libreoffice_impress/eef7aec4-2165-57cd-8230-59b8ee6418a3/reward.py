"""
Reward Script: Two-level bulleted list on slide 2
Task ID: impstruct_030
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Main items use square bullets (U+25A0)
  Component 2 (0.30): Sub-items use en-dash bullets (U+2013)
  Component 3 (0.20): Sub-items are at indent level 1
  Component 4 (0.20): Sub-items have increased left margin vs main items
"""

import os
from pptx import Presentation
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impstruct_030'

# Expected structure on slide 2 (after the title paragraph):
# Index 0: "Platform Migration"       -> main item (square bullet, level 0)
# Index 1: "Database schema redesign..." -> sub-item (en-dash, level 1)
# Index 2: "API endpoint compatibility..." -> sub-item (en-dash, level 1)
# Index 3: "Security Hardening"        -> main item (square bullet, level 0)
# Index 4: "Multi-factor authentication..." -> sub-item (en-dash, level 1)
# Index 5: "Penetration testing..."     -> sub-item (en-dash, level 1)

MAIN_INDICES = [0, 3]       # indices into the bullet list (excluding title)
SUB_INDICES = [1, 2, 4, 5]  # indices into the bullet list (excluding title)

SQUARE_BULLET = '\u25a0'  # U+25A0
EN_DASH = '\u2013'        # U+2013


def get_bullet_paragraphs(slide):
    """Get non-title text paragraphs from slide 2's content shapes."""
    paragraphs = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs.append(para)
    # The first paragraph is the title ("Deliverables"), skip it
    # Return only the bullet items
    if len(paragraphs) > 1:
        return paragraphs[1:]
    return []


def get_buChar(para):
    """Extract the bullet character from a paragraph's XML."""
    pPr = para._p.find(qn('a:pPr'))
    if pPr is not None:
        bc = pPr.find(qn('a:buChar'))
        if bc is not None:
            return bc.get('char')
    return None


def get_marL(para):
    """Extract left margin (marL) from paragraph properties."""
    pPr = para._p.find(qn('a:pPr'))
    if pPr is not None:
        val = pPr.get('marL')
        if val is not None:
            return int(val)
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 2 slides
    if len(prs.slides) < 2:
        print("FAIL: Presentation has fewer than 2 slides")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[1]  # slide 2 (0-indexed)
    bullet_paras = get_bullet_paragraphs(slide)

    if len(bullet_paras) < 6:
        print(f"FAIL: Expected 6 bullet paragraphs on slide 2, found {len(bullet_paras)}")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Found {len(bullet_paras)} bullet paragraphs on slide 2")

    # Component 1: Main items (indices 0,3) use square bullets (0.30 points)
    try:
        main_square_count = 0
        for idx in MAIN_INDICES:
            char = get_buChar(bullet_paras[idx])
            text = bullet_paras[idx].text.strip()
            if char == SQUARE_BULLET:
                main_square_count += 1
                print(f"  PASS: Main item '{text}' has square bullet")
            else:
                print(f"  FAIL: Main item '{text}' expected square bullet '{SQUARE_BULLET}', found {char!r}")
        if main_square_count == len(MAIN_INDICES):
            print(f"PASS: Component 1 - All main items have square bullets (0.30 pts)")
            total_score += 0.30
        elif main_square_count > 0:
            partial = 0.30 * (main_square_count / len(MAIN_INDICES))
            print(f"PARTIAL: Component 1 - {main_square_count}/{len(MAIN_INDICES)} main items have square bullets ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 - No main items have square bullets")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Sub-items (indices 1,2,4,5) use en-dash bullets (0.30 points)
    try:
        sub_dash_count = 0
        for idx in SUB_INDICES:
            char = get_buChar(bullet_paras[idx])
            text = bullet_paras[idx].text.strip()
            if char == EN_DASH:
                sub_dash_count += 1
                print(f"  PASS: Sub-item '{text[:40]}...' has en-dash bullet")
            else:
                print(f"  FAIL: Sub-item '{text[:40]}...' expected en-dash '{EN_DASH}', found {char!r}")
        if sub_dash_count == len(SUB_INDICES):
            print(f"PASS: Component 2 - All sub-items have en-dash bullets (0.30 pts)")
            total_score += 0.30
        elif sub_dash_count > 0:
            partial = 0.30 * (sub_dash_count / len(SUB_INDICES))
            print(f"PARTIAL: Component 2 - {sub_dash_count}/{len(SUB_INDICES)} sub-items have en-dash bullets ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 - No sub-items have en-dash bullets")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Sub-items are at indent level 1 (0.20 points)
    try:
        sub_level_count = 0
        for idx in SUB_INDICES:
            lvl = bullet_paras[idx].level
            text = bullet_paras[idx].text.strip()
            if lvl == 1:
                sub_level_count += 1
                print(f"  PASS: Sub-item '{text[:40]}...' at level 1")
            else:
                print(f"  FAIL: Sub-item '{text[:40]}...' expected level 1, found {lvl}")
        if sub_level_count == len(SUB_INDICES):
            print(f"PASS: Component 3 - All sub-items at indent level 1 (0.20 pts)")
            total_score += 0.20
        elif sub_level_count > 0:
            partial = 0.20 * (sub_level_count / len(SUB_INDICES))
            print(f"PARTIAL: Component 3 - {sub_level_count}/{len(SUB_INDICES)} sub-items at level 1 ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 - No sub-items at indent level 1")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Sub-items have greater left margin than main items (0.20 points)
    try:
        # Get margins for main items
        main_margins = []
        for idx in MAIN_INDICES:
            m = get_marL(bullet_paras[idx])
            main_margins.append(m)
            print(f"  INFO: Main item index {idx} marL={m}")

        # Get margins for sub-items
        sub_margins = []
        for idx in SUB_INDICES:
            m = get_marL(bullet_paras[idx])
            sub_margins.append(m)
            print(f"  INFO: Sub-item index {idx} marL={m}")

        # Check that sub-items have strictly greater margin than main items
        # Use the max main margin as baseline
        main_max = max(m for m in main_margins if m is not None) if any(m is not None for m in main_margins) else 0
        sub_min = min(m for m in sub_margins if m is not None) if any(m is not None for m in sub_margins) else 0

        if sub_min > main_max and sub_min > 0:
            print(f"PASS: Component 4 - Sub-items have increased left margin (sub_min={sub_min} > main_max={main_max}) (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 4 - Sub-items do not have increased left margin (sub_min={sub_min}, main_max={main_max})")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
