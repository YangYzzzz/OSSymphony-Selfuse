"""
Initial Setup: Create a 14-slide Project_Plan presentation with no section dividers.
Task ID: impress_rp_049
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_049'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)

def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide titles and content for a 14-slide project plan presentation
    slides_data = [
        {
            "title": "Project Plan: Digital Transformation Initiative",
            "content": "Q3 2025 - Q1 2026\nPresented by: Strategic Planning Division\nVersion 2.3",
            "layout": 0  # Title Slide
        },
        {
            "title": "Executive Summary",
            "content": "This initiative aims to modernize our core business processes through cloud migration, AI integration, and workflow automation. Expected ROI: 340% over 3 years with a total investment of $2.4M.",
            "layout": 1
        },
        {
            "title": "Project Team & Stakeholders",
            "content": "Project Lead: Sarah Chen\nTechnical Director: Marcus Johnson\nBusiness Analyst: Priya Patel\nUX Lead: David Kim\nQA Manager: Elena Rodriguez\nSponsor: VP of Operations, James Thompson",
            "layout": 1
        },
        {
            "title": "Market Research Findings",
            "content": "• 78% of competitors have adopted cloud infrastructure\n• Customer satisfaction scores dropped 12% due to slow processes\n• Industry benchmark: 99.9% uptime, current: 97.2%\n• Average response time target: <200ms",
            "layout": 1
        },
        {
            "title": "Requirements Gathering",
            "content": "• 47 stakeholder interviews completed\n• 12 process workflows documented\n• 8 critical pain points identified\n• Integration requirements with 5 legacy systems\n• Security compliance: SOC 2 Type II, GDPR",
            "layout": 1
        },
        {
            "title": "Technical Architecture",
            "content": "Frontend: React 18 with TypeScript\nBackend: Python FastAPI microservices\nDatabase: PostgreSQL 15 + Redis cache\nCloud: AWS (ECS, RDS, ElastiCache)\nCI/CD: GitHub Actions + ArgoCD",
            "layout": 1
        },
        {
            "title": "Development Sprints Overview",
            "content": "Sprint 1-3: Core API development\nSprint 4-6: Frontend implementation\nSprint 7-8: Integration testing\nSprint 9: Performance optimization\nSprint 10: Security audit & fixes",
            "layout": 1
        },
        {
            "title": "Sprint Velocity & Metrics",
            "content": "Average velocity: 42 story points/sprint\nBurn-down rate: 93% completion\nDefect density: 0.8 per 1000 LOC\nCode coverage: 87%\nTechnical debt ratio: 4.2%",
            "layout": 1
        },
        {
            "title": "Deployment Strategy",
            "content": "Phase 1: Canary deployment (5% traffic)\nPhase 2: Blue-green rollout (50%)\nPhase 3: Full production cutover\nRollback plan: Automated within 15 minutes\nMonitoring: Datadog + PagerDuty",
            "layout": 1
        },
        {
            "title": "Quality Assurance Results",
            "content": "• 2,847 test cases executed\n• 98.7% pass rate\n• 12 critical bugs found and resolved\n• Load testing: 10,000 concurrent users\n• Accessibility: WCAG 2.1 AA compliant",
            "layout": 1
        },
        {
            "title": "Post-Launch Review",
            "content": "• System uptime: 99.95% (first 30 days)\n• Response time: avg 145ms (target <200ms)\n• User adoption: 89% within first week\n• Support tickets: 23 (down from 156/month)\n• NPS score improved: 34 → 67",
            "layout": 1
        },
        {
            "title": "Lessons Learned",
            "content": "• Early stakeholder engagement reduced scope changes by 60%\n• Automated testing saved 340 hours of manual QA\n• Cloud cost optimization reduced monthly spend by 28%\n• Cross-team communication tools improved delivery speed",
            "layout": 1
        },
        {
            "title": "Budget Summary",
            "content": "Total Budget: $2,400,000\nPersonnel: $1,680,000 (70%)\nInfrastructure: $480,000 (20%)\nLicensing: $168,000 (7%)\nContingency: $72,000 (3%)\nActual Spend: $2,287,500 (4.7% under budget)",
            "layout": 1
        },
        {
            "title": "Next Steps & Recommendations",
            "content": "1. Scale to remaining 3 business units by Q2 2026\n2. Implement ML-based anomaly detection\n3. Expand API marketplace for partners\n4. Quarterly security penetration testing\n5. Annual architecture review cycle",
            "layout": 1
        },
    ]

    for i, sd in enumerate(slides_data):
        layout_idx = sd["layout"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = sd["title"]
            for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                run.font.size = Pt(28)
                run.font.bold = True

        # Set content
        if layout_idx == 0:
            # Title slide - use subtitle placeholder
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = sd["content"]
        elif layout_idx == 1:
            # Title + Content layout
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = sd["content"]

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')

create_initial()
