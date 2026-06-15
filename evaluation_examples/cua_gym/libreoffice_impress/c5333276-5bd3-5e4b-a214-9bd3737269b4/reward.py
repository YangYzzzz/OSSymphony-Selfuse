"""
Reward Script: Section divider slide system verification
Task ID: impress_rp_049
Domain: libreoffice_impress
Scoring:
  Component 1: Slide count is 17 (0.15)
  Component 2: Divider at position 4 with correct bg #1ABC9C (0.15)
  Component 3: Divider at position 9 with correct bg #E74C3C (0.15)
  Component 4: Divider at position 14 with correct bg #9B59B6 (0.15)
  Component 5: Divider 1 text 'Phase 1: Discovery' in 44pt bold white centered (0.15)
  Component 6: Divider 2 text 'Phase 2: Execution' in 44pt bold white centered (0.10)
  Component 7: Divider 3 text 'Phase 3: Review' in 44pt bold white centered (0.15)
"""

import os
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_049'


def get_slide_background_rgb(slide):
    """Extract background color from slide."""
    try:
        fill = slide.background.fill
        if fill.type is not None and fill.type == 1:  # solid fill
            return str(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def get_all_text_shapes(slide):
    """Get all shapes with text frames, including nested in groups."""
    def extract(shape):
        results = []
        if hasattr(shape, "text") and hasattr(shape, "text_frame"):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.extend(extract(sub))
        return results
    out = []
    for shape in slide.shapes:
        out.extend(extract(shape))
    return out


def check_divider_text(slide, expected_text):
    """Check that slide has the expected text with 44pt bold white centered formatting."""
    text_shapes = get_all_text_shapes(slide)

    for shape in text_shapes:
        for para in shape.text_frame.paragraphs:
            full_text = para.text.strip()
            if full_text == expected_text:
                # Check alignment is centered
                alignment_ok = (para.alignment == PP_ALIGN.CENTER)

                # Check runs for font properties
                non_empty_runs = [r for r in para.runs if (r.text or "").strip()]
                if not non_empty_runs:
                    return False, "No non-empty runs found"

                for run in non_empty_runs:
                    # Check font size (44pt = 44 * 12700 EMU)
                    if run.font.size is None or abs(run.font.size - Pt(44)) > Pt(1):
                        return False, f"Font size mismatch: expected 44pt, got {run.font.size}"

                    # Check bold (None means inherit, treat as not bold)
                    if run.font.bold is not True:
                        return False, f"Not bold: {run.font.bold}"

                    # Check white color
                    try:
                        if run.font.color.type is not None:
                            color_str = str(run.font.color.rgb)
                            if color_str.upper() != "FFFFFF":
                                return False, f"Color not white: {color_str}"
                        else:
                            return False, "Font color type is None"
                    except Exception as e:
                        return False, f"Color check error: {e}"

                if not alignment_ok:
                    return False, f"Not centered: alignment={para.alignment}"

                return True, "All checks passed"

    return False, f"Text '{expected_text}' not found in any shape"


def verify_task(file_path):
    """Verify task completion with progressive scoring."""
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Slide count is 17 (0.15 points)
    try:
        if num_slides == 17:
            print(f"PASS: Component 1 - Slide count is 17 (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 - Expected 17 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Define expected dividers (0-indexed positions: 3, 8, 13)
    dividers = [
        {"pos": 3, "bg": "1ABC9C", "text": "Phase 1: Discovery"},
        {"pos": 8, "bg": "E74C3C", "text": "Phase 2: Execution"},
        {"pos": 13, "bg": "9B59B6", "text": "Phase 3: Review"},
    ]

    # Component 2: Divider 1 background #1ABC9C at position 4 (0.15 pts)
    try:
        if num_slides >= 4:
            slide = prs.slides[3]
            bg_color = get_slide_background_rgb(slide)
            if bg_color and bg_color.upper() == "1ABC9C":
                print(f"PASS: Component 2 - Slide 4 background is #1ABC9C (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 2 - Slide 4 background expected #1ABC9C, found {bg_color}")
        else:
            print(f"FAIL: Component 2 - Not enough slides for position 4")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Divider 2 background #E74C3C at position 9 (0.15 pts)
    try:
        if num_slides >= 9:
            slide = prs.slides[8]
            bg_color = get_slide_background_rgb(slide)
            if bg_color and bg_color.upper() == "E74C3C":
                print(f"PASS: Component 3 - Slide 9 background is #E74C3C (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 3 - Slide 9 background expected #E74C3C, found {bg_color}")
        else:
            print(f"FAIL: Component 3 - Not enough slides for position 9")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Divider 3 background #9B59B6 at position 14 (0.15 pts)
    try:
        if num_slides >= 14:
            slide = prs.slides[13]
            bg_color = get_slide_background_rgb(slide)
            if bg_color and bg_color.upper() == "9B59B6":
                print(f"PASS: Component 4 - Slide 14 background is #9B59B6 (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 4 - Slide 14 background expected #9B59B6, found {bg_color}")
        else:
            print(f"FAIL: Component 4 - Not enough slides for position 14")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Divider 1 text 'Phase 1: Discovery' 44pt bold white centered (0.15 pts)
    try:
        if num_slides >= 4:
            slide = prs.slides[3]
            ok, detail = check_divider_text(slide, "Phase 1: Discovery")
            if ok:
                print(f"PASS: Component 5 - Slide 4 text 'Phase 1: Discovery' formatted correctly (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 5 - {detail}")
        else:
            print(f"FAIL: Component 5 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    # Component 6: Divider 2 text 'Phase 2: Execution' 44pt bold white centered (0.10 pts)
    try:
        if num_slides >= 9:
            slide = prs.slides[8]
            ok, detail = check_divider_text(slide, "Phase 2: Execution")
            if ok:
                print(f"PASS: Component 6 - Slide 9 text 'Phase 2: Execution' formatted correctly (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 6 - {detail}")
        else:
            print(f"FAIL: Component 6 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 6 - {e}")

    # Component 7: Divider 3 text 'Phase 3: Review' 44pt bold white centered (0.15 pts)
    try:
        if num_slides >= 14:
            slide = prs.slides[13]
            ok, detail = check_divider_text(slide, "Phase 3: Review")
            if ok:
                print(f"PASS: Component 7 - Slide 14 text 'Phase 3: Review' formatted correctly (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 7 - {detail}")
        else:
            print(f"FAIL: Component 7 - Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 7 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
