"""
Reward Script: Update title text and colors on slides 3 and 5 to match slides 1 and 2
Task ID: osworld_impress_title_color_match_005
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Slide 3 title text == 'Methodology Overview'
  Component 2 (0.25): Slide 3 title color == #1A237E (dark navy, matching slide 1)
  Component 3 (0.25): Slide 5 title text == 'Statistical Analysis Results'
  Component 4 (0.25): Slide 5 title color == #228B22 (forest green, matching slide 2)
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_color_match_005'

# Expected ground truth values from task context
SLIDE3_EXPECTED_TITLE = 'Methodology Overview'
SLIDE3_EXPECTED_COLOR = '1A237E'  # dark navy — matching slide 1

SLIDE5_EXPECTED_TITLE = 'Statistical Analysis Results'
SLIDE5_EXPECTED_COLOR = '228B22'  # forest green — matching slide 2


def get_title_shape(slide):
    """Return the title placeholder (idx==0) for a slide, or None."""
    for shape in slide.shapes:
        if (shape.has_text_frame
                and hasattr(shape, 'placeholder_format')
                and shape.placeholder_format is not None
                and shape.placeholder_format.idx == 0):
            return shape
    return None


def get_title_text(shape):
    """Return the full text of a title shape's first paragraph."""
    if shape is None:
        return None
    try:
        return shape.text_frame.paragraphs[0].text.strip()
    except Exception:
        return None


def get_title_color(shape):
    """
    Return the RGB color string of the first non-empty run in the title,
    or None if not set as an explicit RGB color.
    """
    if shape is None:
        return None
    try:
        para = shape.text_frame.paragraphs[0]
        runs = [r for r in para.runs if (r.text or '').strip()]
        if not runs:
            return None
        run = runs[0]
        if run.font.color.type is not None:
            return str(run.font.color.rgb).upper()
        return None
    except Exception:
        return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition gate: file must be loadable
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition gate: must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"CRITICAL: Expected at least 5 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]  # 0-based index for slide 3
    slide5 = prs.slides[4]  # 0-based index for slide 5

    title3_shape = get_title_shape(slide3)
    title5_shape = get_title_shape(slide5)

    # Component 1: Slide 3 title text updated to 'Methodology Overview' (0.25 points)
    # Initial text was 'Methods'; must be changed to the new title
    try:
        actual_text3 = get_title_text(title3_shape)
        if actual_text3 == SLIDE3_EXPECTED_TITLE:
            print(f"PASS: Component 1 — Slide 3 title is '{actual_text3}' (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — Slide 3 title expected '{SLIDE3_EXPECTED_TITLE}', found '{actual_text3}'")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 3 title color matches slide 1 (#1A237E dark navy) (0.25 points)
    # Initial color was #000000 (black); must be changed to #1A237E
    try:
        actual_color3 = get_title_color(title3_shape)
        if actual_color3 == SLIDE3_EXPECTED_COLOR:
            print(f"PASS: Component 2 — Slide 3 title color is #{actual_color3} (matches slide 1 navy, 0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 — Slide 3 title color expected #{SLIDE3_EXPECTED_COLOR}, found #{actual_color3}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 5 title text updated to 'Statistical Analysis Results' (0.25 points)
    # Initial text was 'Results'; must be changed to the new title
    try:
        actual_text5 = get_title_text(title5_shape)
        if actual_text5 == SLIDE5_EXPECTED_TITLE:
            print(f"PASS: Component 3 — Slide 5 title is '{actual_text5}' (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 — Slide 5 title expected '{SLIDE5_EXPECTED_TITLE}', found '{actual_text5}'")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 5 title color matches slide 2 (#228B22 forest green) (0.25 points)
    # Initial color was #000000 (black); must be changed to #228B22
    try:
        actual_color5 = get_title_color(title5_shape)
        if actual_color5 == SLIDE5_EXPECTED_COLOR:
            print(f"PASS: Component 4 — Slide 5 title color is #{actual_color5} (matches slide 2 green, 0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 4 — Slide 5 title color expected #{SLIDE5_EXPECTED_COLOR}, found #{actual_color5}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in given env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
