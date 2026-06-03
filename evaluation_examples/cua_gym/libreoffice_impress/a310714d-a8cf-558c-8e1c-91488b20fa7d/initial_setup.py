"""
Initial Setup: Research paper presentation with 8 slides.
Slide 3 title reads 'Methods' in black; slide 5 reads 'Results' in black.
The agent must rename and recolor both to match other slides' title colors.
Task ID: osworld_impress_title_color_match_005
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_color_match_005'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_title_color(slide, color_rgb):
    """Set the title placeholder text color on a slide."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.shape_type == 13:
            continue
        if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
            ph_idx = shape.placeholder_format.idx
            if ph_idx == 0:  # title placeholder
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = color_rgb
                return


def set_title_text_and_color(slide, title_text, color_rgb):
    """Set title text and color on a slide."""
    for shape in slide.shapes:
        if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
            ph_idx = shape.placeholder_format.idx
            if ph_idx == 0:  # title placeholder
                tf = shape.text_frame
                tf.text = title_text
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = color_rgb
                        run.font.size = Pt(32)
                        run.font.bold = True
                return


def add_content_text(slide, content_lines):
    """Add content text to the body placeholder of a slide."""
    for shape in slide.shapes:
        if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
            ph_idx = shape.placeholder_format.idx
            if ph_idx == 1:  # body/content placeholder
                tf = shape.text_frame
                tf.text = content_lines[0] if content_lines else ""
                for line in content_lines[1:]:
                    p = tf.add_paragraph()
                    p.text = line
                    p.level = 1
                return


def create_initial():
    prs = Presentation()
    # Use standard widescreen dimensions
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Color palette for slides
    NAVY = RGBColor(0x1A, 0x23, 0x7E)        # slide 1 title: dark navy
    FOREST_GREEN = RGBColor(0x22, 0x8B, 0x22) # slide 2 title: forest green
    BLACK = RGBColor(0x00, 0x00, 0x00)        # slides 3 and 5 (to be updated)
    DARK_TEAL = RGBColor(0x00, 0x69, 0x5C)   # slide 4
    MAROON = RGBColor(0x88, 0x00, 0x36)       # slide 6
    DARK_PURPLE = RGBColor(0x4A, 0x14, 0x8C) # slide 7
    DARK_GRAY = RGBColor(0x37, 0x47, 0x4F)   # slide 8

    # Define slide data: (title, title_color, content_lines)
    slides_data = [
        # Slide 1: Introduction (dark navy)
        (
            "Introduction",
            NAVY,
            [
                "Background and Motivation",
                "Recent advances in computational methods",
                "Gap in current literature identified",
                "Scope and objectives of this study",
            ],
        ),
        # Slide 2: Literature Review (forest green)
        (
            "Literature Review",
            FOREST_GREEN,
            [
                "Prior work on quantitative analysis",
                "Chen et al. (2022): baseline framework",
                "Martinez & Patel (2023): improved accuracy",
                "Key limitations of existing approaches",
            ],
        ),
        # Slide 3: Methods — black, MUST NOT be 'Methodology Overview' or navy
        (
            "Methods",
            BLACK,
            [
                "Study design: cross-sectional cohort",
                "Participants: 240 subjects (ages 25–65)",
                "Data collection instruments",
                "Statistical procedures applied",
            ],
        ),
        # Slide 4: Data Collection (dark teal)
        (
            "Data Collection",
            DARK_TEAL,
            [
                "Survey instruments administered online",
                "Response rate: 78.4% (n=188 valid responses)",
                "Data cleaning and preprocessing steps",
                "Exclusion criteria applied to 12 records",
            ],
        ),
        # Slide 5: Results — black, MUST NOT be 'Statistical Analysis Results' or green
        (
            "Results",
            BLACK,
            [
                "Primary outcome: significant improvement (p<0.001)",
                "Secondary outcomes: 3 of 4 hypotheses confirmed",
                "Effect sizes ranged from 0.42 to 0.71",
                "No adverse interactions detected",
            ],
        ),
        # Slide 6: Discussion (maroon)
        (
            "Discussion",
            MAROON,
            [
                "Findings consistent with Chen et al. (2022)",
                "Unexpected result in sub-group analysis",
                "Practical implications for practitioners",
                "Theoretical contributions to the field",
            ],
        ),
        # Slide 7: Conclusion (dark purple)
        (
            "Conclusion",
            DARK_PURPLE,
            [
                "Study objectives were met",
                "Strong evidence for proposed framework",
                "Recommendations for policy implementation",
                "Future research directions identified",
            ],
        ),
        # Slide 8: References (dark gray)
        (
            "References",
            DARK_GRAY,
            [
                "Chen, L., et al. (2022). Journal of Computational Methods, 14(3), 112-134.",
                "Martinez, R. & Patel, S. (2023). Data Science Review, 7(1), 45-67.",
                "Thompson, A., et al. (2021). Quantitative Analysis Today, 9(2), 88-105.",
                "Williams, J. (2020). Statistical Approaches in Research, 5th ed. Oxford.",
            ],
        ),
    ]

    layout_title_content = prs.slide_layouts[1]  # Title and Content

    for title_text, title_color, content_lines in slides_data:
        slide = prs.slides.add_slide(layout_title_content)
        set_title_text_and_color(slide, title_text, title_color)
        add_content_text(slide, content_lines)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready: open the file in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
