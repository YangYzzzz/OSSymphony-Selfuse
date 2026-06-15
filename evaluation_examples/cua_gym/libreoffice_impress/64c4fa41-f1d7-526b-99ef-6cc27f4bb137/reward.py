"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 213, Picture 1 sits a bit lower than Picture 2. How do I use LibreOffice Impress’s ‘Align → Top’ command to lock both images to the exact same Y-position so their top edges line up perfectly?
Generated: 2025-09-10 19:10:54
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def verify_alignment(file_path: str) -> float:
    """Verify that on slide 213 the first two pictures are perfectly top-aligned.

    Scoring (progressive):
      • 0.3 – Target slide exists in the file
      • 0.3 – At least two picture shapes exist on that slide
      • 0.4 – The two pictures share (within tolerance) the exact same Y (top) position
    Returns a float between 0.0 and 1.0, printing detailed diagnostics.
    """

    MAX_SCORE = 1.0
    score = 0.0
    TARGET_SLIDE_NUM = 213           # 1-based index requested in the task
    ALIGN_TOLERANCE_EMU = 1000       # ≈0.01 pt – practically identical

    # ---------- 1) File loading ----------
    print(f"Verifying presentation file: {file_path}")

    if not os.path.exists(file_path):
        print("✗ File not found")
        return 0.0

    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded successfully with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Could not open presentation: {e}")
        return 0.0

    # ---------- 2) Locate target slide ----------
    if len(prs.slides) < TARGET_SLIDE_NUM:
        print(f"✗ Presentation contains fewer than {TARGET_SLIDE_NUM} slides")
        return 0.0

    slide = prs.slides[TARGET_SLIDE_NUM - 1]
    print(f"✓ Found target slide #{TARGET_SLIDE_NUM}")
    score += 0.3  # earned for correct slide presence

    # ---------- 3) Find picture shapes ----------
    pictures = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    print(f"Found {len(pictures)} picture shape(s) on the slide")

    if len(pictures) < 2:
        print("✗ Need at least two pictures to verify alignment")
        return score  # partial credit only

    score += 0.3  # pictures present

    # Prefer shapes explicitly named "Picture 1" and "Picture 2" if available
    pic1 = None
    pic2 = None
    for pic in pictures:
        lower_name = pic.name.lower().strip()
        if pic1 is None and lower_name == "picture 1":
            pic1 = pic
        elif pic2 is None and lower_name == "picture 2":
            pic2 = pic

    # Fallback to the first two pictures if specific names not found
    if pic1 is None or pic2 is None:
        pic1, pic2 = pictures[0], pictures[1]

    # ---------- 4) Verify top alignment ----------
    top_diff = abs(pic1.top - pic2.top)
    print(f"Top positions – Picture 1: {pic1.top}, Picture 2: {pic2.top}, difference: {top_diff}")

    if top_diff <= ALIGN_TOLERANCE_EMU:
        print("✓ Pictures are top-aligned within tolerance")
        score += 0.4
    else:
        print("✗ Pictures are NOT properly top-aligned")

    final_score = min(score, MAX_SCORE)
    print(f"Total Score: {final_score}/{MAX_SCORE}")
    return final_score


if __name__ == "__main__":
    FILE_PATH = "/home/user/on_slide_213_picture_1_sits_a_bit_lower_than_picture_2_how_do_i_use_libreoffice_impresss_align_top_c_golden.pptx"
    reward = verify_alignment(FILE_PATH)
    print(f"REWARD: {reward}")
