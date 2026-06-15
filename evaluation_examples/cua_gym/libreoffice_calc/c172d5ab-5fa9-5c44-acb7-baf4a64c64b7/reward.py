"""
Reward Script: Customer Satisfaction Presentation from PDF Data
Task ID: pdf_cross_064
Domain: libreoffice_impress + pdf (cross-domain)
Scoring:
  Component 1: ODP file exists with 7 slides                              (0.25 pts)
  Component 2: Each of 5 category slides has progress bar shapes          (0.35 pts)
  Component 3: Summary slide references all 5 categories                  (0.20 pts)
  Component 4: Exported PDF exists with 7 pages                           (0.20 pts)
"""

import os

# Domain-specific imports
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import pymupdf
    PYMUPDF_AVAILABLE = True
except ImportError:
    try:
        import fitz as pymupdf
        PYMUPDF_AVAILABLE = True
    except ImportError:
        PYMUPDF_AVAILABLE = False

WORKDIR = '/home/user/Documents'
ODP_PATH = f'{WORKDIR}/satisfaction_report.odp'
PDF_PATH = f'{WORKDIR}/satisfaction_report.pdf'

# Ground truth from task context
EXPECTED_SLIDE_COUNT = 7
EXPECTED_PDF_PAGES = 7
EXPECTED_CATEGORIES = [
    'Product Quality',
    'Customer Service',
    'Delivery Speed',
    'Value for Money',
    'Overall Experience',
]
# Category slides are slides 2-6 (index 1-5)
CATEGORY_SLIDE_INDICES = [1, 2, 3, 4, 5]
# Summary slide is slide 7 (index 6)
SUMMARY_SLIDE_INDEX = 6


def verify_task():
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # -------------------------------------------------------------------------
    # Component 1: ODP file exists and has 7 slides (0.25 points)
    # This FAILS on initial (file does not exist) and PASSES on golden.
    # -------------------------------------------------------------------------
    try:
        if not PPTX_AVAILABLE:
            print("FAIL: Component 1 — python-pptx not available on this VM")
        elif not os.path.exists(ODP_PATH):
            print(f"FAIL: Component 1 — ODP file not found at {ODP_PATH}")
        else:
            prs = Presentation(ODP_PATH)
            actual_slides = len(prs.slides)
            if actual_slides == EXPECTED_SLIDE_COUNT:
                print(f"PASS: Component 1 — ODP exists with {actual_slides} slides (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 — Expected {EXPECTED_SLIDE_COUNT} slides, found {actual_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # -------------------------------------------------------------------------
    # Component 2: Each of 5 category slides contains visual progress bar shapes
    # A progress bar is represented by AUTO_SHAPE (type=1) rectangles.
    # We require at least 2 AUTO_SHAPE rectangles per category slide
    # (one for the background bar, one for the filled/coloured portion).
    # Scoring: 0.07 pts per category slide, 5 slides × 0.07 = 0.35 pts total.
    # This FAILS on initial (ODP not present) and PASSES on golden.
    # -------------------------------------------------------------------------
    try:
        if not PPTX_AVAILABLE:
            print("FAIL: Component 2 — python-pptx not available on this VM")
        elif not os.path.exists(ODP_PATH):
            print("FAIL: Component 2 — ODP file not found, cannot check progress bars")
        else:
            prs = Presentation(ODP_PATH)
            slide_count = len(prs.slides)
            category_pass = 0
            for slide_idx in CATEGORY_SLIDE_INDICES:
                if slide_idx >= slide_count:
                    print(f"FAIL: Component 2 — Slide {slide_idx + 1} does not exist")
                    continue
                slide = prs.slides[slide_idx]
                # Count AUTO_SHAPE rectangles (shape_type=1) — progress bar shapes
                auto_shapes = [s for s in slide.shapes if s.shape_type == 1]
                if len(auto_shapes) >= 2:
                    print(f"PASS: Component 2 — Slide {slide_idx + 1} has {len(auto_shapes)} progress bar shapes (0.07 pts)")
                    total_score += 0.07
                    category_pass += 1
                else:
                    print(f"FAIL: Component 2 — Slide {slide_idx + 1} has only {len(auto_shapes)} AUTO_SHAPE(s), need >= 2 for progress bar")
            print(f"Component 2 summary: {category_pass}/5 category slides have progress bars")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # -------------------------------------------------------------------------
    # Component 3: Summary slide (slide 7) references all 5 categories (0.20 points)
    # The summary slide must contain text mentioning all 5 category names.
    # This FAILS on initial (ODP not present) and PASSES on golden.
    # -------------------------------------------------------------------------
    try:
        if not PPTX_AVAILABLE:
            print("FAIL: Component 3 — python-pptx not available on this VM")
        elif not os.path.exists(ODP_PATH):
            print("FAIL: Component 3 — ODP file not found, cannot check summary slide")
        else:
            prs = Presentation(ODP_PATH)
            slide_count = len(prs.slides)
            if SUMMARY_SLIDE_INDEX >= slide_count:
                print(f"FAIL: Component 3 — Summary slide (index {SUMMARY_SLIDE_INDEX}) does not exist")
            else:
                summary_slide = prs.slides[SUMMARY_SLIDE_INDEX]
                # Collect all text from summary slide shapes
                all_text_parts = []
                for shape in summary_slide.shapes:
                    if shape.has_text_frame:
                        all_text_parts.append(shape.text_frame.text)
                full_text = ' '.join(all_text_parts)

                # Check that all 5 categories are mentioned
                missing = [cat for cat in EXPECTED_CATEGORIES if cat not in full_text]
                if not missing:
                    print(f"PASS: Component 3 — Summary slide references all 5 categories (0.20 pts)")
                    total_score += 0.20
                else:
                    print(f"FAIL: Component 3 — Summary slide missing categories: {missing}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # -------------------------------------------------------------------------
    # Component 4: Exported PDF exists and has 7 pages (0.20 points)
    # This FAILS on initial (PDF not present) and PASSES on golden.
    # -------------------------------------------------------------------------
    try:
        if not PYMUPDF_AVAILABLE:
            print("FAIL: Component 4 — pymupdf not available, cannot verify PDF")
        elif not os.path.exists(PDF_PATH):
            print(f"FAIL: Component 4 — PDF file not found at {PDF_PATH}")
        else:
            doc = pymupdf.open(PDF_PATH)
            actual_pages = doc.page_count
            doc.close()
            if actual_pages == EXPECTED_PDF_PAGES:
                print(f"PASS: Component 4 — PDF exists with {actual_pages} pages (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 4 — Expected {EXPECTED_PDF_PAGES} pages, found {actual_pages}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Cap at 1.0 to handle floating-point accumulation (5 × 0.07 = 0.35)
    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


verify_task()
