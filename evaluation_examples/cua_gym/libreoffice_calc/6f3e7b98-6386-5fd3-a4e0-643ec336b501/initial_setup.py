"""
Initial Setup: Batch convert .ppt and .pptx files to PDF using terminal commands
Task ID: osworld_multi_apps_batch_convert_015
Domain: multi_apps (LibreOffice Impress + Terminal)
"""

import os
import shlex
import subprocess
import time

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_batch_convert_015'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_pptx_file(filepath, title, subtitle, slides_data):
    """Create a .pptx file using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = subtitle

    # Content slides
    content_layout = prs.slide_layouts[1]
    for slide_title, slide_content in slides_data:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = slide_title
        tf = slide.placeholders[1].text_frame
        tf.text = slide_content

    prs.save(filepath)
    print(f'Created: {filepath}')


def create_ppt_stub(filepath, title, content_lines):
    """
    Create a minimal .ppt (PowerPoint 97-2003) file using python-pptx saved as .ppt
    by writing it as a pptx first then renaming, or use subprocess with libreoffice to convert.
    Since python-pptx cannot write .ppt natively, create a .pptx then convert to .ppt
    using libreoffice --headless.
    """
    # First create a temp pptx
    tmp_pptx = filepath.replace('.ppt', '_tmp.pptx')
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    # Title slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = content_lines[0] if content_lines else ''

    # Additional content slides
    content_layout = prs.slide_layouts[1]
    for i, line in enumerate(content_lines[1:], 1):
        s = prs.slides.add_slide(content_layout)
        s.shapes.title.text = f'Slide {i + 1}'
        s.placeholders[1].text = line

    prs.save(tmp_pptx)

    # Convert to .ppt using libreoffice
    outdir = os.path.dirname(filepath)
    result = subprocess.run(
        ['libreoffice', '--headless', '--convert-to', 'ppt', '--outdir', outdir, tmp_pptx],
        capture_output=True, text=True, timeout=60
    )
    # The output file from libreoffice conversion will be named based on the input
    converted = tmp_pptx.replace('_tmp.pptx', '_tmp.ppt')
    if os.path.exists(converted):
        os.rename(converted, filepath)
        print(f'Converted to .ppt: {filepath}')
    else:
        # Fallback: just rename pptx to ppt (not ideal but workable for setup)
        print(f'WARNING: .ppt conversion failed, stdout={result.stdout}, stderr={result.stderr}')
        print(f'Falling back to copying pptx as ppt: {filepath}')
        import shutil
        shutil.copy(tmp_pptx, filepath)

    # Clean up temp file
    if os.path.exists(tmp_pptx):
        os.remove(tmp_pptx)


def create_initial():
    os.makedirs(DESKTOP, exist_ok=True)

    # Remove any pre-existing PDF files on Desktop (ensure clean initial state)
    for f in os.listdir(DESKTOP):
        if f.endswith('.pdf'):
            os.remove(os.path.join(DESKTOP, f))
            print(f'Removed pre-existing PDF: {f}')

    # --- Create .pptx files ---
    pptx_files = [
        {
            'filename': 'Q3_Sales_Review.pptx',
            'title': 'Q3 Sales Review 2024',
            'subtitle': 'Regional Performance & Forecasts',
            'slides': [
                ('Revenue Summary', 'Total Q3 revenue reached $4.2M, up 18% from Q2.\nNorth America led with $2.1M in sales.'),
                ('Top Products', 'Product A: $1.4M\nProduct B: $980K\nProduct C: $720K'),
                ('Action Items', 'Expand Product A into APAC markets\nHire 3 additional sales reps in EMEA\nLaunch Q4 promotions by Oct 1'),
            ]
        },
        {
            'filename': 'Marketing_Strategy_2025.pptx',
            'title': 'Marketing Strategy 2025',
            'subtitle': 'Brand Growth & Digital Initiatives',
            'slides': [
                ('Executive Summary', 'Focus on digital-first campaigns across social media platforms.\nTarget 25% increase in brand awareness.'),
                ('Campaign Roadmap', 'Q1: Brand refresh and website redesign\nQ2: Social media blitz\nQ3: Influencer partnerships\nQ4: Holiday season campaigns'),
                ('Budget Allocation', 'Digital advertising: 45%\nContent creation: 25%\nEvents & sponsorships: 20%\nReserve: 10%'),
            ]
        },
        {
            'filename': 'HR_Onboarding_Training.pptx',
            'title': 'Employee Onboarding Program',
            'subtitle': 'Welcome to the Team',
            'slides': [
                ('Company Overview', 'Founded in 2010, we serve 5,000+ enterprise clients worldwide.\nHeadquartered in San Francisco with offices in London and Singapore.'),
                ('Benefits & Perks', 'Health, dental, and vision insurance\n401(k) with 4% company match\nFlexible working hours and remote options\n20 days PTO plus 10 company holidays'),
                ('First Week Schedule', 'Day 1: IT setup and team introductions\nDay 2-3: Product training\nDay 4: Shadow senior colleagues\nDay 5: Wrap-up and feedback session'),
            ]
        },
    ]

    for info in pptx_files:
        filepath = os.path.join(DESKTOP, info['filename'])
        create_pptx_file(filepath, info['title'], info['subtitle'], info['slides'])

    # --- Create .ppt files ---
    ppt_files = [
        {
            'filename': 'Project_Kickoff_2024.ppt',
            'title': 'Project Kickoff: Atlas Initiative',
            'lines': [
                'Scope: Build a unified data platform for enterprise analytics.',
                'Timeline: 6-month delivery in 3 phases.',
                'Budget: $850,000 total project investment.',
            ]
        },
        {
            'filename': 'Annual_Report_FY2023.ppt',
            'title': 'Annual Report FY2023',
            'lines': [
                'Revenue: $28.4M (+12% YoY)\nGross Margin: 63%\nNet Income: $5.1M',
                'Key Milestones: Launched 2 new product lines, expanded to 3 new markets.',
                'Outlook FY2024: Target $33M revenue with focus on subscription growth.',
            ]
        },
        {
            'filename': 'Tech_Architecture_Overview.ppt',
            'title': 'Technology Architecture Overview',
            'lines': [
                'Current stack: Python microservices, PostgreSQL, Redis, AWS.',
                'Proposed improvements: Migrate to Kubernetes, adopt GraphQL API layer.',
                'Security: Zero-trust model with MFA and end-to-end encryption.',
            ]
        },
    ]

    for info in ppt_files:
        filepath = os.path.join(DESKTOP, info['filename'])
        create_ppt_stub(filepath, info['title'], info['lines'])

    # Verify all files are present
    desktop_files = os.listdir(DESKTOP)
    pptx_count = sum(1 for f in desktop_files if f.endswith('.pptx'))
    ppt_count = sum(1 for f in desktop_files if f.endswith('.ppt'))
    pdf_count = sum(1 for f in desktop_files if f.endswith('.pdf'))
    print(f'\nDesktop contents: {pptx_count} .pptx files, {ppt_count} .ppt files, {pdf_count} .pdf files')
    print('Expected: 3 .pptx, 3 .ppt, 0 .pdf')

    # GUI-ready startup: open terminal so the agent can run commands
    launch_gui('gnome-terminal', delay_sec=2.0)
    print('GUI_READY: launched gnome-terminal with DISPLAY=:0')


create_initial()
