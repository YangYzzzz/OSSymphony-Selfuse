"""
Reward Script: Multi-app workflow — Extract PDF invoices, create Calc spreadsheet,
create Impress presentation, export to PDF.
Task ID: pdf_cross_142
Domain: libreoffice_calc (multi-app: calc + impress + pdf)
Scoring:
  Component 1: XLSX has 4 sheets (Vendor_A, Vendor_B, Vendor_C, Summary)          — 0.30
  Component 2: XLSX summary values correct (A=$3200, B=$5750, C=$2100, GT=$11050) — 0.25
  Component 3: PPTX has 4 slides with expected slide titles/content               — 0.20
  Component 4: invoice_review.pdf exists with 4 pages                             — 0.15
  Component 5: PDF contains key financial data (totals, vendor names)             — 0.10
  Total: 1.0
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'pdf_cross_142'

XLSX_PATH = f'{WORKDIR}/{TASK_ID}.xlsx'
PPTX_PATH = f'{WORKDIR}/{TASK_ID}.pptx'
PDF_PATH = f'{WORKDIR}/Documents/invoice_review.pdf'

# Expected ground truth values from task context
EXPECTED_VENDOR_A_TOTAL = 3200
EXPECTED_VENDOR_B_TOTAL = 5750
EXPECTED_VENDOR_C_TOTAL = 2100
EXPECTED_GRAND_TOTAL = 11050


def _check_xlsx_sheets():
    """Check XLSX has 4 expected sheets. Returns score (0..0.30)."""
    import openpyxl
    if not os.path.exists(XLSX_PATH):
        print(f"FAIL: Component 1 — XLSX file not found: {XLSX_PATH}")
        return 0.0
    wb = openpyxl.load_workbook(XLSX_PATH)
    sheet_names = wb.sheetnames
    expected_sheets = ['Vendor_A', 'Vendor_B', 'Vendor_C', 'Summary']
    sheets_present = [s for s in expected_sheets if s in sheet_names]
    sheets_score = len(sheets_present) / len(expected_sheets)
    if sheets_score == 1.0:
        print(f"PASS: Component 1 — XLSX has all 4 required sheets: {sheet_names} (0.30 pts)")
        return 0.30
    elif sheets_score > 0:
        partial = round(0.30 * sheets_score, 4)
        print(f"PARTIAL: Component 1 — XLSX has {len(sheets_present)}/4 sheets: "
              f"found {sheet_names}, expected {expected_sheets} ({partial:.2f} pts)")
        return partial
    else:
        print(f"FAIL: Component 1 — XLSX exists but lacks expected sheets: "
              f"found {sheet_names}, expected {expected_sheets}")
        return 0.0


def _check_xlsx_values():
    """Check XLSX Summary sheet totals. Returns score (0..0.25)."""
    import openpyxl
    if not os.path.exists(XLSX_PATH):
        print("FAIL: Component 2 — XLSX file not found")
        return 0.0
    wb = openpyxl.load_workbook(XLSX_PATH, data_only=True)
    if 'Summary' not in wb.sheetnames:
        print("FAIL: Component 2 — Summary sheet not found in XLSX")
        return 0.0
    ws_summary = wb['Summary']

    found_a = False
    found_b = False
    found_c = False
    found_grand = False

    for row in ws_summary.iter_rows():
        for cell in row:
            val = cell.value
            if val is None:
                continue
            try:
                num = float(val)
                if abs(num - EXPECTED_VENDOR_A_TOTAL) < 0.01:
                    found_a = True  # inside conditional: abs(num - 3200) < 0.01
                if abs(num - EXPECTED_VENDOR_B_TOTAL) < 0.01:
                    found_b = True  # inside conditional: abs(num - 5750) < 0.01
                if abs(num - EXPECTED_VENDOR_C_TOTAL) < 0.01:
                    found_c = True  # inside conditional: abs(num - 2100) < 0.01
                if abs(num - EXPECTED_GRAND_TOTAL) < 0.01:
                    found_grand = True  # inside conditional: abs(num - 11050) < 0.01
            except (ValueError, TypeError):
                pass

    checks = [found_a, found_b, found_c, found_grand]
    labels = [
        f"Vendor A total={EXPECTED_VENDOR_A_TOTAL}",
        f"Vendor B total={EXPECTED_VENDOR_B_TOTAL}",
        f"Vendor C total={EXPECTED_VENDOR_C_TOTAL}",
        f"Grand total={EXPECTED_GRAND_TOTAL}",
    ]
    passed = sum(checks)
    sub_score = round(0.25 * (passed / len(checks)), 4)

    for label, ok in zip(labels, checks):
        status = "PASS" if ok else "FAIL"
        print(f"  {status}: Summary {label}")

    if passed == len(checks):
        print(f"PASS: Component 2 — All 4 summary totals correct (0.25 pts)")
    elif passed > 0:
        print(f"PARTIAL: Component 2 — {passed}/4 summary totals correct ({sub_score:.2f} pts)")
    else:
        print(f"FAIL: Component 2 — No summary totals found in Summary sheet")

    return sub_score


def _check_pptx_slides():
    """Check PPTX has 4 slides with expected content. Returns score (0..0.20)."""
    from pptx import Presentation
    if not os.path.exists(PPTX_PATH):
        print(f"FAIL: Component 3 — PPTX file not found: {PPTX_PATH}")
        return 0.0

    prs = Presentation(PPTX_PATH)
    num_slides = len(prs.slides)

    if num_slides != 4:
        print(f"FAIL: Component 3 — Expected 4 slides, found {num_slides}")
        # Partial: file exists but wrong slide count
        if num_slides > 0:
            return round(0.20 * min(num_slides / 4.0, 1.0) * 0.5, 4)
        return 0.0

    sub_score_3 = 0.0

    # Sub-check 3a: Slide 1 has title text about invoices/review/report
    slide1_texts = []
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    slide1_texts.append(t.lower())

    slide1_has_title = any(
        'invoice' in t or 'review' in t or 'report' in t
        for t in slide1_texts
    )
    if slide1_has_title:
        print("PASS: Component 3a — Slide 1 has invoice/review title text")
        sub_score_3 += 0.05
    else:
        print(f"FAIL: Component 3a — Slide 1 lacks expected title text, found: {slide1_texts[:3]}")

    # Sub-check 3b: Slide 2 has vendor spending figures ($3,200, $5,750, $2,100)
    slide2_texts = []
    for shape in prs.slides[1].shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    slide2_texts.append(t)

    slide2_all_text = ' '.join(slide2_texts)
    spending_found = sum(1 for amt in ['3,200', '5,750', '2,100']
                         if amt in slide2_all_text or amt.replace(',', '') in slide2_all_text)
    if spending_found >= 2:
        print(f"PASS: Component 3b — Slide 2 has vendor spending data ({spending_found}/3 amounts)")
        sub_score_3 += 0.05
    else:
        print(f"FAIL: Component 3b — Slide 2 missing spending amounts; found {spending_found}/3")

    # Sub-check 3c: Slide 3 has table with at least 2 rows (header + 1+ expense rows)
    slide3_has_table = False
    for shape in prs.slides[2].shapes:
        if shape.shape_type == 19:  # TABLE
            table = shape.table
            if len(table.rows) >= 2:
                slide3_has_table = True
                if slide3_has_table:
                    print(f"PASS: Component 3c — Slide 3 has expense table ({len(table.rows)} rows)")
                    sub_score_3 += 0.05
                    break

    if not slide3_has_table:
        print("FAIL: Component 3c — Slide 3 lacks expense table with data rows")

    # Sub-check 3d: Slide 4 has summary content (table or summary text)
    slide4_texts = []
    slide4_has_table = False
    for shape in prs.slides[3].shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    slide4_texts.append(t.lower())
        if shape.shape_type == 19:  # TABLE
            slide4_has_table = True

    slide4_all = ' '.join(slide4_texts)
    slide4_ok = (
        slide4_has_table or
        'summary' in slide4_all or
        'grand' in slide4_all or
        '11,050' in slide4_all or
        '11050' in slide4_all
    )
    if slide4_ok:
        print("PASS: Component 3d — Slide 4 has summary content")
        sub_score_3 += 0.05
    else:
        print(f"FAIL: Component 3d — Slide 4 lacks summary content; found: {slide4_texts[:5]}")

    print(f"Component 3 total: {sub_score_3:.2f}/0.20 pts")
    return sub_score_3


def _check_pdf_pages():
    """Check invoice_review.pdf has 4 pages. Returns score (0..0.15)."""
    if not os.path.exists(PDF_PATH):
        print(f"FAIL: Component 4 — invoice_review.pdf not found at {PDF_PATH}")
        return 0.0
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(PDF_PATH)
        page_count = doc.page_count
        doc.close()
        if page_count == 4:
            print(f"PASS: Component 4 — invoice_review.pdf has 4 pages (0.15 pts)")
            return 0.15
        elif page_count > 0:
            partial = round(0.15 * min(page_count / 4.0, 1.0), 4)
            print(f"PARTIAL: Component 4 — invoice_review.pdf has {page_count} pages "
                  f"(expected 4) ({partial:.2f} pts)")
            return partial
        else:
            print(f"FAIL: Component 4 — invoice_review.pdf is empty (0 pages)")
            return 0.0
    except ImportError:
        # fitz not available — use file size as a proxy for a non-trivial PDF
        file_size = os.path.getsize(PDF_PATH)
        if file_size > 1000:
            print(f"PARTIAL: Component 4 — invoice_review.pdf exists ({file_size} bytes); "
                  f"fitz unavailable for page count (0.07 pts)")
            return 0.07
        else:
            print(f"FAIL: Component 4 — invoice_review.pdf is too small ({file_size} bytes)")
            return 0.0


def _check_pdf_content():
    """Check invoice_review.pdf contains key financial data. Returns score (0..0.10)."""
    if not os.path.exists(PDF_PATH):
        print("FAIL: Component 5 — invoice_review.pdf not found")
        return 0.0
    try:
        import fitz
        doc = fitz.open(PDF_PATH)
        all_text = ""
        for page in doc:
            all_text += page.get_text("text")
        doc.close()
        all_text_lower = all_text.lower()

        checks_5 = [
            ('11,050' in all_text or '11050' in all_text,
             "Grand total $11,050"),
            ('techsupplies' in all_text_lower or 'vendor a' in all_text_lower,
             "Vendor A (TechSupplies)"),
            ('officegear' in all_text_lower or 'officesupply' in all_text_lower,
             "Vendor B or C (OfficeGear/OfficeSupply)"),
        ]

        passed_5 = sum(1 for ok, _ in checks_5 if ok)
        for ok, label in checks_5:
            status = "PASS" if ok else "FAIL"
            print(f"  {status}: Component 5 — PDF contains {label}")

        sub_score_5 = round(0.10 * (passed_5 / len(checks_5)), 4)
        if passed_5 == len(checks_5):
            print(f"PASS: Component 5 — All key financial data found in PDF (0.10 pts)")
        elif passed_5 > 0:
            print(f"PARTIAL: Component 5 — {passed_5}/{len(checks_5)} checks passed "
                  f"({sub_score_5:.2f} pts)")
        else:
            print(f"FAIL: Component 5 — PDF lacks expected financial data")

        return sub_score_5

    except ImportError:
        print("SKIP: Component 5 — fitz not available, skipping PDF content check")
        return 0.0


def verify_task():
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Component 1: XLSX has 4 correct sheets (0.30 points)
    try:
        score_1 = _check_xlsx_sheets()
        if score_1 > 0:
            total_score += score_1
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: XLSX Summary totals correct (0.25 points)
    try:
        score_2 = _check_xlsx_values()
        if score_2 > 0:
            total_score += score_2
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: PPTX has 4 slides with expected content (0.20 points)
    try:
        score_3 = _check_pptx_slides()
        if score_3 > 0:
            total_score += score_3
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: invoice_review.pdf has 4 pages (0.15 points)
    try:
        score_4 = _check_pdf_pages()
        if score_4 > 0:
            total_score += score_4
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: PDF contains key financial data (0.10 points)
    try:
        score_5 = _check_pdf_content()
        if score_5 > 0:
            total_score += score_5
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(round(total_score, 4), 1.0)
    print(f"\nScore: {total_score:.4f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


if __name__ == '__main__':
    verify_task()
