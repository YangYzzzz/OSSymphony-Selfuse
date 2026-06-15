"""
Reward Script: Extract frames from film_clip.mp4 and set as backgrounds in Cinema_Analysis.pptx
Task ID: osworld_multi_apps_vlc_frame_to_slide_008
Domain: multi_apps (VLC + LibreOffice Impress)

Scoring Rubric:
  Component 1: Three frame PNG files exist on Desktop (0.3 pts)
               - frame_05.png (frame at 00:05)
               - frame_30.png (frame at 00:30)
               - frame_60.png (frame at 01:00)
  Component 2: Slides 1, 2, 3 each have a full-slide PICTURE shape as background (0.4 pts)
               - Each picture covers the full slide (left=0, top=0, same width/height as slide)
  Component 3: The picture in each slide matches the corresponding frame PNG file (0.3 pts)
               - Slide 1 picture blob == frame_05.png content
               - Slide 2 picture blob == frame_30.png content
               - Slide 3 picture blob == frame_60.png content

Total: 1.0
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_vlc_frame_to_slide_008'
PPTX_PATH = f'{WORKDIR}/Cinema_Analysis.pptx'
FRAME_FILES = {
    1: 'frame_05.png',
    2: 'frame_30.png',
    3: 'frame_60.png',
}


def verify_task():
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: PPTX file must exist
    if not os.path.exists(PPTX_PATH):
        print(f"CRITICAL: Presentation file not found: {PPTX_PATH}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Load the PPTX
    try:
        prs = Presentation(PPTX_PATH)
    except Exception as e:
        print(f"CRITICAL: Cannot load presentation {PPTX_PATH}: {e}")
        print("REWARD: 0.0")
        return 0.0

    if len(prs.slides) < 3:
        print(f"CRITICAL: Expected at least 3 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # -----------------------------------------------------------------------
    # Component 1: Three frame PNG files exist on Desktop (0.3 pts)
    # These files should NOT exist in initial_env, only in golden_env
    # -----------------------------------------------------------------------
    try:
        frames_found = []
        for slide_idx, fname in FRAME_FILES.items():
            fpath = os.path.join(WORKDIR, fname)
            if os.path.exists(fpath) and os.path.getsize(fpath) > 0:
                frames_found.append(fname)
                print(f"PASS: Frame file exists: {fpath} ({os.path.getsize(fpath)} bytes)")
            else:
                print(f"FAIL: Frame file missing or empty: {fpath}")

        if len(frames_found) == 3:
            print(f"PASS: Component 1 — All 3 frame PNG files found on Desktop (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 — Only {len(frames_found)}/3 frame files found ({frames_found})")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # -----------------------------------------------------------------------
    # Component 2: Slides 1, 2, 3 each have a full-slide PICTURE shape (0.4 pts)
    # The picture must span the entire slide (left=0, top=0, full width and height)
    # This is absent in initial_env
    # -----------------------------------------------------------------------
    try:
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        full_slide_pics = {}

        for slide_num in [1, 2, 3]:
            slide = prs.slides[slide_num - 1]
            found_full_pic = False
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Check if it covers the full slide
                    if (shape.left == 0 and shape.top == 0 and
                            shape.width == slide_width and shape.height == slide_height):
                        found_full_pic = True
                        print(f"PASS: Slide {slide_num} has a full-slide PICTURE shape "
                              f"({shape.width}x{shape.height} at 0,0)")
                        break
            full_slide_pics[slide_num] = found_full_pic
            if not found_full_pic:
                print(f"FAIL: Slide {slide_num} does NOT have a full-slide PICTURE shape")

        if all(full_slide_pics.values()):
            print(f"PASS: Component 2 — All 3 slides have full-slide picture backgrounds (0.4 pts)")
            total_score += 0.4
        else:
            missing = [s for s, v in full_slide_pics.items() if not v]
            print(f"FAIL: Component 2 — Slides {missing} lack full-slide picture backgrounds")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # -----------------------------------------------------------------------
    # Component 3: Picture content in each slide matches corresponding frame file (0.3 pts)
    # Slide 1 picture blob == frame_05.png, etc.
    # This requires both the files AND the matching blobs to be present
    # -----------------------------------------------------------------------
    try:
        blob_matches = {}

        for slide_num, fname in FRAME_FILES.items():
            fpath = os.path.join(WORKDIR, fname)
            if not os.path.exists(fpath):
                print(f"FAIL: Component 3 — Frame file missing for comparison: {fpath}")
                blob_matches[slide_num] = False
                continue

            with open(fpath, 'rb') as f:
                frame_data = f.read()

            slide = prs.slides[slide_num - 1]
            matched = False
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    blob = shape.image.blob
                    if frame_data == blob:
                        matched = True
                        print(f"PASS: Slide {slide_num} picture blob matches {fname} exactly")
                        break
            if not matched:
                print(f"FAIL: Slide {slide_num} picture blob does NOT match {fname}")
            blob_matches[slide_num] = matched

        if all(blob_matches.values()):
            print(f"PASS: Component 3 — All 3 slide pictures match their corresponding frame files (0.3 pts)")
            total_score += 0.3
        else:
            mismatched = [s for s, v in blob_matches.items() if not v]
            print(f"FAIL: Component 3 — Slides {mismatched} picture content does not match frames")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


if __name__ == '__main__':
    verify_task()
