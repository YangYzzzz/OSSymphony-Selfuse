"""
Initial Setup: VLC frame extraction and LibreOffice Impress background task
Task ID: osworld_multi_apps_vlc_frame_to_slide_008
Domain: multi_apps (VLC + LibreOffice Impress)

Creates:
  - /home/user/Desktop/film_clip.mp4  : a test video clip with distinct colored scenes
  - /home/user/Desktop/Cinema_Analysis.pptx : 5-slide presentation with white backgrounds on slides 1-3
  - Does NOT create frame_05.png, frame_30.png, frame_60.png (agent must extract these)
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_vlc_frame_to_slide_008'
PPTX_PATH = f'{DESKTOP}/Cinema_Analysis.pptx'
VIDEO_PATH = f'{DESKTOP}/film_clip.mp4'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_video():
    """Generate a test MP4 video with distinct colored scenes at different timestamps."""
    # Ensure Desktop exists
    os.makedirs(DESKTOP, exist_ok=True)

    # Create a 75-second video with distinct colored segments using ffmpeg
    # Segment at 0-15s: red scene (so frame at 00:05 is red)
    # Segment at 15-45s: green scene (so frame at 00:30 is green)
    # Segment at 45-75s: blue scene (so frame at 01:00 is blue)
    # Each segment has unique color so frames can be distinguished
    filter_complex = (
        "color=c=red:s=1280x720:r=25:d=15[seg1];"
        "color=c=green:s=1280x720:r=25:d=30[seg2];"
        "color=c=blue:s=1280x720:r=25:d=30[seg3];"
        "[seg1][seg2][seg3]concat=n=3:v=1:a=0[out]"
    )

    result = subprocess.run([
        "ffmpeg", "-y",
        "-filter_complex", filter_complex,
        "-map", "[out]",
        "-pix_fmt", "yuv420p",
        "-c:v", "libx264",
        VIDEO_PATH
    ], capture_output=True, text=True)

    if result.returncode != 0:
        print(f"ffmpeg error: {result.stderr}")
        # Fallback: simpler single color video if complex filter fails
        subprocess.run([
            "ffmpeg", "-y",
            "-f", "lavfi",
            "-i", "color=c=darkblue:s=1280x720:r=25:d=75",
            "-pix_fmt", "yuv420p",
            "-c:v", "libx264",
            VIDEO_PATH
        ], check=True)

    print(f"Video created: {VIDEO_PATH}")


def create_presentation():
    """Create Cinema_Analysis.pptx with 5 slides, slides 1-3 have white backgrounds."""
    prs = Presentation()

    # Set standard widescreen dimensions
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_layout_blank = prs.slide_layouts[6]   # Blank layout
    slide_layout_title = prs.slide_layouts[0]   # Title Slide
    slide_layout_content = prs.slide_layouts[1] # Title and Content

    # --- Slide 1: Title slide with white background ---
    slide1 = prs.slides.add_slide(slide_layout_blank)
    # Set explicit white background
    fill1 = slide1.background.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Add title text box
    txb1 = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(2))
    tf1 = txb1.text_frame
    p1 = tf1.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = "Cinema Analysis"
    run1.font.size = Pt(48)
    run1.font.bold = True
    run1.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

    p1b = tf1.add_paragraph()
    p1b.alignment = PP_ALIGN.CENTER
    run1b = p1b.add_run()
    run1b.text = "Film Study: Visual Composition and Cinematography"
    run1b.font.size = Pt(20)
    run1b.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # --- Slide 2: Opening sequence with white background ---
    slide2 = prs.slides.add_slide(slide_layout_blank)
    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    txb2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1.2))
    tf2 = txb2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "Act I — Opening Sequence (00:05)"
    run2.font.size = Pt(32)
    run2.font.bold = True
    run2.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    txb2b = slide2.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.33), Inches(3))
    tf2b = txb2b.text_frame
    tf2b.word_wrap = True
    p2b = tf2b.paragraphs[0]
    run2b = p2b.add_run()
    run2b.text = (
        "The opening frames establish the visual tone and narrative context. "
        "Key cinematographic techniques observed include the use of natural lighting "
        "and wide establishing shots to orient the viewer within the story's world."
    )
    run2b.font.size = Pt(16)
    run2b.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 3: Middle act with white background ---
    slide3 = prs.slides.add_slide(slide_layout_blank)
    fill3 = slide3.background.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    txb3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1.2))
    tf3 = txb3.text_frame
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = "Act II — Rising Action (00:30)"
    run3.font.size = Pt(32)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    txb3b = slide3.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.33), Inches(3))
    tf3b = txb3b.text_frame
    tf3b.word_wrap = True
    p3b = tf3b.paragraphs[0]
    run3b = p3b.add_run()
    run3b.text = (
        "The middle act demonstrates increasing tension through deliberate camera work. "
        "Close-up shots and shallow depth of field draw attention to character emotions, "
        "while the color grading shifts toward warmer tones to reflect narrative escalation."
    )
    run3b.font.size = Pt(16)
    run3b.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 4: Climax analysis (dark background, non-white) ---
    slide4 = prs.slides.add_slide(slide_layout_blank)
    fill4 = slide4.background.fill
    fill4.solid()
    fill4.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    txb4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1.2))
    tf4 = txb4.text_frame
    p4 = tf4.paragraphs[0]
    run4 = p4.add_run()
    run4.text = "Act III — Climax (01:00)"
    run4.font.size = Pt(32)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    txb4b = slide4.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.33), Inches(3))
    tf4b = txb4b.text_frame
    tf4b.word_wrap = True
    p4b = tf4b.paragraphs[0]
    run4b = p4b.add_run()
    run4b.text = (
        "The climactic sequence employs rapid editing and dynamic camera movements. "
        "High contrast lighting and desaturated color palettes create a sense of urgency. "
        "This section represents the film's peak dramatic intensity."
    )
    run4b.font.size = Pt(16)
    run4b.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    # --- Slide 5: Conclusion and Summary ---
    slide5 = prs.slides.add_slide(slide_layout_blank)
    fill5 = slide5.background.fill
    fill5.solid()
    fill5.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)

    txb5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1.2))
    tf5 = txb5.text_frame
    p5 = tf5.paragraphs[0]
    run5 = p5.add_run()
    run5.text = "Conclusions & Cinematographic Assessment"
    run5.font.size = Pt(32)
    run5.font.bold = True
    run5.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    txb5b = slide5.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.33), Inches(4))
    tf5b = txb5b.text_frame
    tf5b.word_wrap = True
    for line in [
        "Director demonstrates mastery of visual storytelling through consistent motifs",
        "Color grading evolution mirrors character arcs effectively",
        "Camera movement vocabulary is purposeful and contextually appropriate",
        "Editing rhythm maintains audience engagement throughout the narrative",
        "Overall cinematographic quality: Exceptional",
    ]:
        p_item = tf5b.add_paragraph()
        run_item = p_item.add_run()
        run_item.text = f"• {line}"
        run_item.font.size = Pt(15)
        run_item.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(PPTX_PATH)
    print(f"Presentation created: {PPTX_PATH}")


def verify_no_frames():
    """Ensure frame PNG files do not exist (agent must create them)."""
    frame_files = [
        f'{DESKTOP}/frame_05.png',
        f'{DESKTOP}/frame_30.png',
        f'{DESKTOP}/frame_60.png',
    ]
    for f in frame_files:
        if os.path.exists(f):
            os.remove(f)
            print(f"Removed pre-existing frame: {f}")


def main():
    os.makedirs(DESKTOP, exist_ok=True)

    # 1. Create the video clip
    create_video()

    # 2. Create the presentation
    create_presentation()

    # 3. Ensure frame PNGs don't exist
    verify_no_frames()

    # 4. GUI-ready startup: open VLC with the video and LibreOffice Impress with the presentation
    # Launch LibreOffice Impress first
    launch_gui(f'libreoffice --impress "{PPTX_PATH}"', delay_sec=3.0)

    # Launch VLC with the video file
    launch_gui(f'vlc "{VIDEO_PATH}"', delay_sec=2.0)

    print("GUI_READY: launched LibreOffice Impress and VLC with DISPLAY=:0")


main()
