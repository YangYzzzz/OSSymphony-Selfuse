"""
Initial Setup: VLC playing film score video + LibreOffice Impress with cinema_club.pptx
Task ID: osworld_multi_apps_misc_065
Domain: multi_apps (VLC + LibreOffice Impress)
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_misc_065'
PPTX_PATH = f'{WORKDIR}/cinema_club.pptx'
MP4_PATH = f'{DESKTOP}/film_score_compilation.mp4'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env['DISPLAY'] = ':0'
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_mp4():
    """Create a short film score compilation video on the Desktop using ffmpeg."""
    os.makedirs(DESKTOP, exist_ok=True)

    # Generate a short video (10 seconds) with audio - film score style
    # Use ffmpeg to create a test video with sine wave audio at 440Hz
    result = subprocess.run([
        'ffmpeg', '-y',
        '-f', 'lavfi',
        '-i', 'testsrc=duration=10:size=1280x720:rate=24',
        '-f', 'lavfi',
        '-i', 'sine=frequency=440:duration=10',
        '-c:v', 'libx264',
        '-c:a', 'aac',
        '-shortest',
        MP4_PATH
    ], capture_output=True, text=True)

    if result.returncode != 0:
        # Try alternate approach with lavfi overlay
        result = subprocess.run([
            'ffmpeg', '-y',
            '-f', 'lavfi',
            '-i', 'color=c=black:size=640x480:duration=10:rate=24',
            '-f', 'lavfi',
            '-i', 'anullsrc=r=44100:cl=stereo',
            '-t', '10',
            '-c:v', 'libx264',
            '-c:a', 'aac',
            MP4_PATH
        ], capture_output=True, text=True)

    if result.returncode == 0:
        print(f'MP4 created: {MP4_PATH}')
    else:
        print(f'ffmpeg error: {result.stderr}')
        raise RuntimeError('Failed to create MP4')


def create_pptx():
    """Create cinema_club.pptx with 4 slides in /home/user."""
    prs = Presentation()

    # Slide 1: Title Slide - Cinema Club Welcome
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = 'Cinema Club'
    slide1.placeholders[1].text = 'An Evening of Cinematic Excellence'

    # Style title
    title_frame = slide1.shapes.title.text_frame
    for para in title_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(44)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # Background
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF5, 0xF0, 0xE8)

    # Slide 2: This Month's Feature Films
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "This Month's Feature Films"

    tf2 = slide2.placeholders[1].text_frame
    tf2.text = 'The Grand Budapest Hotel'
    films = [
        'Blade Runner 2049',
        'Interstellar',
        'La La Land',
        'The Shawshank Redemption',
    ]
    for film in films:
        p = tf2.add_paragraph()
        p.text = film
        p.level = 1

    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    for para in slide2.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xFF, 0xD7, 0x00)

    # Slide 3: Upcoming Screenings
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = 'Upcoming Screenings'

    tf3 = slide3.placeholders[1].text_frame
    screenings = [
        'March 15 — Casablanca (1942)',
        'March 22 — 2001: A Space Odyssey (1968)',
        'March 29 — Citizen Kane (1941)',
        'April 5 — Vertigo (1958)',
        'April 12 — Sunset Boulevard (1950)',
    ]
    tf3.text = screenings[0]
    for item in screenings[1:]:
        p = tf3.add_paragraph()
        p.text = item
        p.level = 0

    fill3 = slide3.background.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(0x2C, 0x2C, 0x54)
    for para in slide3.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Slide 4: Membership & Contact
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = 'Join the Cinema Club'

    tf4 = slide4.placeholders[1].text_frame
    tf4.text = 'Annual Membership: $45'
    details = [
        'Monthly screenings every Friday at 7:30 PM',
        'Post-screening discussion and refreshments',
        'Access to our 500+ film library',
        'Contact: cinemaclub@community.org',
        'Website: www.cinemaclub.org',
    ]
    for detail in details:
        p = tf4.add_paragraph()
        p.text = detail
        p.level = 1

    fill4 = slide4.background.fill
    fill4.solid()
    fill4.fore_color.rgb = RGBColor(0x8B, 0x0A, 0x0A)
    for para in slide4.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.bold = True

    prs.save(PPTX_PATH)
    print(f'PPTX created: {PPTX_PATH}')


def create_initial():
    os.makedirs(DESKTOP, exist_ok=True)

    # 1. Create the MP4 video on Desktop
    create_mp4()

    # 2. Create cinema_club.pptx with 4 slides
    create_pptx()

    # 3. Kill any existing VLC/LibreOffice instances for clean state
    subprocess.run(['pkill', '-f', 'vlc'], capture_output=True)
    subprocess.run(['pkill', '-f', 'soffice'], capture_output=True)
    time.sleep(2)

    # 4. Launch VLC playing the film score compilation video
    env = os.environ.copy()
    env['DISPLAY'] = ':0'
    env['VLC_VERBOSE'] = '-1'
    subprocess.Popen(
        ['vlc', MP4_PATH, '--loop'],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(2)
    print('VLC launched with film_score_compilation.mp4')

    # 5. Launch LibreOffice Impress with cinema_club.pptx
    launch_gui(f'libreoffice --impress "{PPTX_PATH}"', delay_sec=3.0)
    print('LibreOffice Impress launched with cinema_club.pptx')

    print('GUI_READY: launched VLC and LibreOffice Impress with DISPLAY=:0')


create_initial()
