"""
Initial Setup: Capture VLC frame at 00:15 and set as slide 1 background
Task ID: osworld_multi_apps_vlc_frame_to_slide_001
Domain: multi_apps (VLC + LibreOffice Impress)

Creates:
  - /home/user/Desktop/savanna_wildlife.mp4  — a savanna wildlife video (>= 20s)
  - /home/user/Wildlife_Deck.pptx            — 5-slide presentation, slide 1 plain white background
Then opens both in VLC and LibreOffice Impress for the GUI agent.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
VIDEO_PATH = f'{DESKTOP}/savanna_wildlife.mp4'
PPTX_PATH = f'{WORKDIR}/Wildlife_Deck.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch a GUI app on the VM display without blocking script exit."""
    env = os.environ.copy()
    env['DISPLAY'] = ':0'
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_video():
    """Generate a synthetic savanna wildlife video (30s, 1280x720) using ffmpeg."""
    os.makedirs(DESKTOP, exist_ok=True)

    # Build a visually interesting video with warm savanna tones using ffmpeg lavfi filters
    # segmented color changes to simulate savanna landscape
    cmd = [
        'ffmpeg', '-y',
        '-f', 'lavfi',
        '-i', (
            'color=c=0x8B6914:size=1280x720:rate=30:duration=30,'
            'format=yuv420p'
        ),
        '-f', 'lavfi',
        '-i', 'anullsrc=r=44100:cl=stereo',
        '-map', '0:v', '-map', '1:a',
        '-c:v', 'libx264', '-pix_fmt', 'yuv420p',
        '-c:a', 'aac', '-shortest',
        '-t', '30',
        VIDEO_PATH,
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        # Fallback: simpler color test video
        cmd2 = [
            'ffmpeg', '-y',
            '-f', 'lavfi', '-i', 'testsrc=duration=30:size=1280x720:rate=30',
            '-f', 'lavfi', '-i', 'anullsrc=r=44100:cl=stereo',
            '-map', '0:v', '-map', '1:a',
            '-c:v', 'libx264', '-pix_fmt', 'yuv420p',
            '-c:a', 'aac', '-shortest',
            '-t', '30',
            VIDEO_PATH,
        ]
        subprocess.run(cmd2, capture_output=True, check=True)
    print(f'Video created: {VIDEO_PATH}')


def create_presentation():
    """Create Wildlife_Deck.pptx with 5 slides; slide 1 has plain white background."""
    prs = Presentation()

    # Slide dimensions (standard 16:9 widescreen)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_data = [
        {
            'layout': 0,  # Title Slide
            'title': 'Savanna Wildlife',
            'subtitle': 'A Journey Through the African Wilderness',
        },
        {
            'layout': 1,  # Title + Content
            'title': 'The Elephant Herd',
            'body': (
                'African elephants are the largest land animals on Earth.\n'
                'Herds are led by the oldest female — the matriarch.\n'
                'They travel up to 50 km per day in search of food and water.\n'
                'A herd can consist of 8–100 individuals.'
            ),
        },
        {
            'layout': 1,
            'title': 'Lions at Dusk',
            'body': (
                'Lions are the only truly social big cats.\n'
                'A pride typically includes 3–30 individuals.\n'
                'The Serengeti supports over 3,000 lions.\n'
                'Females do most of the hunting.'
            ),
        },
        {
            'layout': 1,
            'title': 'The Great Migration',
            'body': (
                'Over 1.5 million wildebeest migrate annually.\n'
                'The circuit spans Tanzania and Kenya.\n'
                'Migration is driven by rainfall and fresh grass.\n'
                'The Mara River crossing is the most dramatic stage.'
            ),
        },
        {
            'layout': 1,
            'title': 'Conservation Efforts',
            'body': (
                'Habitat loss threatens 60% of savanna species.\n'
                'Anti-poaching patrols protect elephant corridors.\n'
                'Community wildlife conservancies cover 12 million hectares.\n'
                'Eco-tourism funds 40% of conservation budgets.'
            ),
        },
    ]

    for i, sd in enumerate(slide_data):
        layout_idx = sd['layout']
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Slide 1: ensure plain white background (no fill = white default)
        if i == 0:
            # Explicitly set white solid background
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = sd['title']
            # Set subtitle
            try:
                slide.placeholders[1].text = sd['subtitle']
            except (KeyError, IndexError):
                pass
        else:
            # Set white background for all other slides too
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = sd['title']
            # Set body/content
            try:
                slide.placeholders[1].text = sd['body']
            except (KeyError, IndexError):
                pass

    prs.save(PPTX_PATH)
    print(f'Presentation created: {PPTX_PATH}')


def main():
    create_video()
    create_presentation()

    # GUI startup: open savanna_wildlife.mp4 in VLC
    launch_gui(f'vlc "{VIDEO_PATH}"', delay_sec=3.0)

    # GUI startup: open Wildlife_Deck.pptx in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{PPTX_PATH}"', delay_sec=3.0)

    print('GUI_READY: launched VLC with savanna_wildlife.mp4 and LibreOffice Impress with Wildlife_Deck.pptx')


main()
