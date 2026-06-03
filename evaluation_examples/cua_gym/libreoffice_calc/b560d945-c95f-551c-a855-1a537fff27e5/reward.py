"""
Reward Script: Capture VLC frame at 00:15 and set as background image of slide 1
Task ID: osworld_multi_apps_vlc_frame_to_slide_001
Domain: libreoffice_impress + vlc (multi-app)
Scoring:
  Component 1: Slide 1 has a PICTURE shape (0.5 pts)
  Component 2: Picture covers the full slide as background (positioned at 0,0 with full width/height) (0.3 pts)
  Component 3: Picture content matches the video frame at ~00:15 (image similarity) (0.2 pts)
  Total: 1.0
"""

import os
import subprocess
import tempfile
import io

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'osworld_multi_apps_vlc_frame_to_slide_001'

PPTX_PATH = f'{WORKDIR}/Wildlife_Deck.pptx'
VIDEO_PATH = f'{WORKDIR}/Desktop/savanna_wildlife.mp4'
TIMESTAMP = '00:00:15'  # 15 seconds


def verify_task(pptx_path, video_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load presentation {pptx_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Verify slide count as a precondition (should have 5 slides)
    if len(prs.slides) < 1:
        print(f"CRITICAL: Presentation has no slides")
        print("REWARD: 0.0")
        return 0.0

    slide1 = prs.slides[0]
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Collect all picture shapes from slide 1
    picture_shapes = [
        s for s in slide1.shapes
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]

    # Component 1: Slide 1 has at least one PICTURE shape (0.5 points)
    # Task requires a frame image to be placed on slide 1.
    # Initial env has 0 picture shapes on slide 1.
    try:
        if len(picture_shapes) >= 1:
            print(f"PASS: Component 1 — Slide 1 has {len(picture_shapes)} picture shape(s) (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 1 — Expected at least 1 picture on slide 1, found 0")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Picture covers the full slide (positioned at 0,0 with full width/height) (0.3 points)
    # Initial env has no picture at all, so this check fails on initial env.
    try:
        if len(picture_shapes) >= 1:
            # Find a picture that covers the full slide (background image)
            covering_pics = [
                pic for pic in picture_shapes
                if (pic.left == 0 and pic.top == 0
                    and abs(pic.width - slide_width) / max(abs(slide_width), 1) <= 0.005
                    and abs(pic.height - slide_height) / max(abs(slide_height), 1) <= 0.005)
            ]
            if len(covering_pics) >= 1:
                cover_pic = covering_pics[0]
                print(f"PASS: Component 2 — Picture covers full slide "
                      f"(left=0, top=0, width={cover_pic.width}, height={cover_pic.height}) (0.3 pts)")
                total_score += 0.3
            else:
                pic = picture_shapes[0]
                print(f"FAIL: Component 2 — Picture does not cover full slide: "
                      f"left={pic.left} (expected 0), top={pic.top} (expected 0), "
                      f"width={pic.width} (expected {slide_width}), "
                      f"height={pic.height} (expected {slide_height})")
        else:
            print(f"FAIL: Component 2 — No picture shape found on slide 1 (skipped)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Picture content matches the video frame at ~00:15 (0.2 points)
    # Uses PIL + numpy for robust image comparison (MSE similarity).
    # Initial env has no picture, so this check fails on initial env.
    try:
        if len(picture_shapes) >= 1 and os.path.exists(video_path):
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_f:
                ref_frame_path = tmp_f.name

            try:
                # Extract reference frame from video at 00:15 using ffmpeg
                result = subprocess.run(
                    ['ffmpeg', '-y', '-i', video_path, '-ss', TIMESTAMP, '-frames:v', '1', ref_frame_path],
                    capture_output=True, timeout=30
                )
                if result.returncode != 0 or not os.path.exists(ref_frame_path) or os.path.getsize(ref_frame_path) == 0:
                    print(f"WARN: Component 3 — Could not extract reference frame from video, skipping")
                else:
                    from PIL import Image
                    import numpy as np

                    # Load reference frame from video
                    ref_img = Image.open(ref_frame_path).convert('RGB')
                    ref_arr = np.array(ref_img, dtype=float)

                    # Load picture from pptx blob
                    pic = picture_shapes[0]
                    pptx_img = Image.open(io.BytesIO(pic.image.blob)).convert('RGB')
                    pptx_arr = np.array(pptx_img, dtype=float)

                    # Resize to same dimensions if needed
                    if ref_img.size != pptx_img.size:
                        smaller_size = (min(ref_img.width, pptx_img.width),
                                        min(ref_img.height, pptx_img.height))
                        ref_img_r = ref_img.resize(smaller_size, Image.Resampling.LANCZOS)
                        pptx_img_r = pptx_img.resize(smaller_size, Image.Resampling.LANCZOS)
                        ref_arr = np.array(ref_img_r, dtype=float)
                        pptx_arr = np.array(pptx_img_r, dtype=float)

                    # Compute normalized mean absolute error
                    mean_abs_diff = np.mean(np.abs(pptx_arr - ref_arr))
                    similarity = 1.0 - (mean_abs_diff / 255.0)
                    print(f"INFO: Component 3 — Image similarity (normalized MAE): {similarity:.4f}, "
                          f"mean_abs_diff={mean_abs_diff:.4f}")

                    # High threshold: content should closely match
                    if similarity >= 0.95:
                        print(f"PASS: Component 3 — Picture content matches video frame at 00:15 "
                              f"(similarity={similarity:.4f}) (0.2 pts)")
                        total_score += 0.2
                    else:
                        print(f"FAIL: Component 3 — Picture content does NOT match video frame at 00:15 "
                              f"(similarity={similarity:.4f}, required >= 0.95)")
            finally:
                if os.path.exists(ref_frame_path):
                    os.unlink(ref_frame_path)

        elif len(picture_shapes) == 0:
            print(f"FAIL: Component 3 — No picture found on slide 1, skipping image comparison")
        else:
            print(f"WARN: Component 3 — Video file not found at {video_path}, skipping image comparison")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


if not os.path.exists(PPTX_PATH):
    print(f"File not found: {PPTX_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(PPTX_PATH, VIDEO_PATH)
