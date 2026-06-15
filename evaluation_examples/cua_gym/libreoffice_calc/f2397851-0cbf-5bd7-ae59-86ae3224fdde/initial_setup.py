"""
Initial Setup: VLC frame extraction and slide background task
Task ID: osworld_multi_apps_vlc_frame_to_slide_007
Domain: multi_apps (VLC + LibreOffice Impress)

Creates:
  - /home/user/Desktop/mountain_sunrise.mp4  (>35s video for frame extraction at 00:30)
  - /home/user/Desktop/Inspiration_Keynote.pptx  (10-slide motivational keynote, slides 1&6 have plain backgrounds)

Opens LibreOffice Impress with the PPTX file.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_vlc_frame_to_slide_007'
VIDEO_PATH = f'{DESKTOP}/mountain_sunrise.mp4'
PPTX_PATH = f'{DESKTOP}/Inspiration_Keynote.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_video():
    """Create a synthetic mountain sunrise video using ffmpeg (>=35 seconds)."""
    os.makedirs(DESKTOP, exist_ok=True)

    if os.path.exists(VIDEO_PATH):
        print(f'Video already exists: {VIDEO_PATH}')
        return

    # Create a sunrise-like gradient video using ffmpeg lavfi
    # Warm orange/golden color to simulate sunrise, 40 seconds long
    cmd = [
        'ffmpeg', '-y',
        '-f', 'lavfi',
        '-i', (
            'color=c=0x1a0a00:size=1280x720:rate=25,'
            'geq='
            "r='clip(255*(t/20),0,255)':"
            "g='clip(120*(t/20),0,180)':"
            "b='clip(30*(t/40),0,80)'"
            ':size=1280x720'
        ),
        '-t', '40',
        '-pix_fmt', 'yuv420p',
        '-vcodec', 'libx264',
        '-crf', '28',
        VIDEO_PATH
    ]

    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        # Fallback: simpler color video
        cmd2 = [
            'ffmpeg', '-y',
            '-f', 'lavfi',
            '-i', 'color=c=orange:size=1280x720:rate=25',
            '-t', '40',
            '-pix_fmt', 'yuv420p',
            VIDEO_PATH
        ]
        result2 = subprocess.run(cmd2, capture_output=True, text=True)
        if result2.returncode != 0:
            # Simplest fallback: testsrc
            cmd3 = [
                'ffmpeg', '-y',
                '-f', 'lavfi',
                '-i', 'testsrc=duration=40:size=1280x720:rate=25',
                '-pix_fmt', 'yuv420p',
                VIDEO_PATH
            ]
            subprocess.run(cmd3, check=True)

    print(f'Video created: {VIDEO_PATH}')


def create_presentation():
    """Create a 10-slide motivational keynote with slides 1 and 6 as plain-background title/section slides."""
    os.makedirs(DESKTOP, exist_ok=True)

    prs = Presentation()

    # Slide dimensions: standard widescreen 13.33" x 7.5"
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Color palette for the keynote
    DARK_NAVY    = RGBColor(0x1A, 0x23, 0x40)   # dark navy blue
    WARM_ORANGE  = RGBColor(0xF2, 0x7D, 0x0A)   # warm orange accent
    LIGHT_GRAY   = RGBColor(0xF4, 0xF4, 0xF4)   # light background
    WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
    DEEP_TEAL    = RGBColor(0x0B, 0x6E, 0x6E)   # section color
    CHARCOAL     = RGBColor(0x2C, 0x2C, 0x2C)

    # Helper: blank slide layout
    blank_layout = prs.slide_layouts[6]   # blank
    title_layout = prs.slide_layouts[0]   # title slide
    content_layout = prs.slide_layouts[1] # title + content

    def add_text_box(slide, text, left, top, width, height,
                     font_name='Calibri', font_size=24,
                     bold=False, color=WHITE, align=PP_ALIGN.LEFT):
        txb = slide.shapes.add_textbox(left, top, width, height)
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        return txb

    def set_bg_solid(slide, rgb):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb

    # ------------------------------------------------------------------
    # SLIDE 1 — Title/Cover Slide (plain dark navy background, no image)
    # ------------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_bg_solid(s1, DARK_NAVY)

    add_text_box(
        s1, 'RISING TO NEW HEIGHTS',
        left=Inches(1.5), top=Inches(2.5), width=Inches(10), height=Inches(1.5),
        font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER
    )
    add_text_box(
        s1, 'A Motivational Keynote on Achievement & Perseverance',
        left=Inches(2.0), top=Inches(4.2), width=Inches(9.0), height=Inches(0.8),
        font_size=20, bold=False, color=WARM_ORANGE, align=PP_ALIGN.CENTER
    )
    add_text_box(
        s1, 'Presented by Jordan Rivera | March 2025',
        left=Inches(3.5), top=Inches(6.0), width=Inches(6.0), height=Inches(0.6),
        font_size=14, bold=False, color=RGBColor(0xBB, 0xBB, 0xBB), align=PP_ALIGN.CENTER
    )

    # ------------------------------------------------------------------
    # SLIDE 2 — Agenda / Overview
    # ------------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    set_bg_solid(s2, LIGHT_GRAY)

    add_text_box(
        s2, 'Today\'s Journey',
        left=Inches(0.8), top=Inches(0.6), width=Inches(9.0), height=Inches(0.9),
        font_size=36, bold=True, color=DARK_NAVY, align=PP_ALIGN.LEFT
    )
    agenda_items = [
        '1.  The Power of Vision',
        '2.  Overcoming Obstacles',
        '3.  Building Resilience',
        '4.  Teamwork & Collaboration',
        '5.  Reaching Your Summit',
    ]
    for i, item in enumerate(agenda_items):
        add_text_box(
            s2, item,
            left=Inches(1.2), top=Inches(1.8 + i * 0.9), width=Inches(10.0), height=Inches(0.75),
            font_size=22, bold=False, color=CHARCOAL, align=PP_ALIGN.LEFT
        )

    # ------------------------------------------------------------------
    # SLIDE 3 — The Power of Vision
    # ------------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    set_bg_solid(s3, DARK_NAVY)

    add_text_box(
        s3, 'THE POWER OF VISION',
        left=Inches(0.8), top=Inches(0.7), width=Inches(11.5), height=Inches(1.0),
        font_size=40, bold=True, color=WARM_ORANGE, align=PP_ALIGN.LEFT
    )
    add_text_box(
        s3, (
            '"Every great dream begins with a dreamer. Always remember, you have '
            'within you the strength, the patience, and the passion to reach for '
            'the stars to change the world."\n\n— Harriet Tubman'
        ),
        left=Inches(1.0), top=Inches(2.0), width=Inches(10.5), height=Inches(3.5),
        font_size=22, bold=False, color=WHITE, align=PP_ALIGN.LEFT
    )
    add_text_box(
        s3, 'Key Insight: Vision without action is a daydream. Action without vision is a nightmare.',
        left=Inches(1.0), top=Inches(5.8), width=Inches(10.5), height=Inches(0.8),
        font_size=16, bold=True, color=WARM_ORANGE, align=PP_ALIGN.LEFT
    )

    # ------------------------------------------------------------------
    # SLIDE 4 — Overcoming Obstacles
    # ------------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    set_bg_solid(s4, LIGHT_GRAY)

    add_text_box(
        s4, 'Overcoming Obstacles',
        left=Inches(0.8), top=Inches(0.6), width=Inches(11.0), height=Inches(1.0),
        font_size=38, bold=True, color=DARK_NAVY, align=PP_ALIGN.LEFT
    )
    obstacles = [
        ('Fear of Failure',    'Reframe failure as data — every setback teaches resilience.'),
        ('Imposter Syndrome',  'Your unique perspective is your greatest asset.'),
        ('Resource Scarcity',  'Constraints fuel creativity. History proves this.'),
        ('Self-Doubt',         'Replace "I can\'t" with "How might I?" and watch doors open.'),
    ]
    for i, (title, desc) in enumerate(obstacles):
        add_text_box(
            s4, f'• {title}',
            left=Inches(1.0), top=Inches(1.8 + i * 1.1), width=Inches(5.5), height=Inches(0.6),
            font_size=18, bold=True, color=DEEP_TEAL, align=PP_ALIGN.LEFT
        )
        add_text_box(
            s4, desc,
            left=Inches(1.2), top=Inches(2.3 + i * 1.1), width=Inches(11.0), height=Inches(0.5),
            font_size=15, bold=False, color=CHARCOAL, align=PP_ALIGN.LEFT
        )

    # ------------------------------------------------------------------
    # SLIDE 5 — Building Resilience
    # ------------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    set_bg_solid(s5, DARK_NAVY)

    add_text_box(
        s5, 'Building Resilience',
        left=Inches(0.8), top=Inches(0.7), width=Inches(11.0), height=Inches(1.0),
        font_size=38, bold=True, color=WARM_ORANGE, align=PP_ALIGN.LEFT
    )
    resilience_points = [
        'Develop a growth mindset — embrace challenges as opportunities',
        'Build your support network before you need it',
        'Practice self-compassion and celebrate small wins',
        'Establish recovery rituals: rest, reflect, re-engage',
        'Document your wins to build evidence of your capability',
    ]
    for i, pt in enumerate(resilience_points):
        add_text_box(
            s5, f'  {i+1}.  {pt}',
            left=Inches(0.8), top=Inches(1.9 + i * 0.95), width=Inches(11.5), height=Inches(0.75),
            font_size=20, bold=False, color=WHITE, align=PP_ALIGN.LEFT
        )

    # ------------------------------------------------------------------
    # SLIDE 6 — Section Divider: "Reaching Your Summit" (plain teal, no image)
    # ------------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    set_bg_solid(s6, DEEP_TEAL)

    add_text_box(
        s6, 'PART II',
        left=Inches(1.5), top=Inches(2.0), width=Inches(10.0), height=Inches(0.8),
        font_size=28, bold=False, color=RGBColor(0xCC, 0xFF, 0xF2), align=PP_ALIGN.CENTER
    )
    add_text_box(
        s6, 'Reaching Your Summit',
        left=Inches(1.0), top=Inches(3.0), width=Inches(11.0), height=Inches(1.5),
        font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER
    )
    add_text_box(
        s6, '"The summit is what drives us, but the climb itself is what matters." — Conrad Anker',
        left=Inches(2.0), top=Inches(5.2), width=Inches(9.0), height=Inches(0.9),
        font_size=16, bold=False, color=RGBColor(0xCC, 0xFF, 0xF2), align=PP_ALIGN.CENTER
    )

    # ------------------------------------------------------------------
    # SLIDE 7 — Teamwork & Collaboration
    # ------------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    set_bg_solid(s7, LIGHT_GRAY)

    add_text_box(
        s7, 'Teamwork & Collaboration',
        left=Inches(0.8), top=Inches(0.6), width=Inches(11.0), height=Inches(1.0),
        font_size=38, bold=True, color=DARK_NAVY, align=PP_ALIGN.LEFT
    )
    team_data = [
        ('Shared Vision',       'Align on WHY before discussing HOW.'),
        ('Psychological Safety', 'Create space where ideas can fail safely.'),
        ('Diverse Perspectives', 'Different backgrounds solve problems faster.'),
        ('Accountability Loop',  'Celebrate wins publicly; address issues privately.'),
    ]
    for i, (heading, body) in enumerate(team_data):
        add_text_box(
            s7, heading,
            left=Inches(0.8), top=Inches(1.9 + i * 1.1), width=Inches(4.5), height=Inches(0.55),
            font_size=18, bold=True, color=DEEP_TEAL, align=PP_ALIGN.LEFT
        )
        add_text_box(
            s7, body,
            left=Inches(5.5), top=Inches(1.9 + i * 1.1), width=Inches(7.5), height=Inches(0.55),
            font_size=17, bold=False, color=CHARCOAL, align=PP_ALIGN.LEFT
        )

    # ------------------------------------------------------------------
    # SLIDE 8 — Case Study: Mount Everest Expedition
    # ------------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    set_bg_solid(s8, DARK_NAVY)

    add_text_box(
        s8, 'Case Study: The Everest Mindset',
        left=Inches(0.8), top=Inches(0.6), width=Inches(11.5), height=Inches(1.0),
        font_size=34, bold=True, color=WARM_ORANGE, align=PP_ALIGN.LEFT
    )
    case_text = (
        'In 1953, Hillary and Tenzing succeeded where 15 expeditions before them failed. '
        'Their secret? They studied every failure, adapted their strategy, and built a '
        'support system of over 400 people.\n\n'
        'Lessons:\n'
        '  • Preparation > Talent\n'
        '  • Learn from every failed attempt\n'
        '  • The right team multiplies individual capability\n'
        '  • Conditions matter — time your push wisely'
    )
    add_text_box(
        s8, case_text,
        left=Inches(1.0), top=Inches(1.8), width=Inches(11.0), height=Inches(5.0),
        font_size=19, bold=False, color=WHITE, align=PP_ALIGN.LEFT
    )

    # ------------------------------------------------------------------
    # SLIDE 9 — Action Plan
    # ------------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    set_bg_solid(s9, LIGHT_GRAY)

    add_text_box(
        s9, 'Your 30-Day Action Plan',
        left=Inches(0.8), top=Inches(0.6), width=Inches(11.0), height=Inches(0.9),
        font_size=36, bold=True, color=DARK_NAVY, align=PP_ALIGN.LEFT
    )
    actions = [
        ('Days 1–5',   'Define your summit: write your 90-day goal in one sentence'),
        ('Days 6–10',  'Map the obstacles: identify top 3 risks and mitigation steps'),
        ('Days 11–18', 'Build your team: schedule 1:1s with 5 key allies'),
        ('Days 19–25', 'Take visible action: complete one high-impact task publicly'),
        ('Days 26–30', 'Review & adapt: measure progress; recalibrate your plan'),
    ]
    for i, (phase, action) in enumerate(actions):
        add_text_box(
            s9, phase,
            left=Inches(0.8), top=Inches(1.75 + i * 1.0), width=Inches(2.5), height=Inches(0.55),
            font_size=16, bold=True, color=WARM_ORANGE, align=PP_ALIGN.LEFT
        )
        add_text_box(
            s9, action,
            left=Inches(3.5), top=Inches(1.75 + i * 1.0), width=Inches(9.5), height=Inches(0.55),
            font_size=16, bold=False, color=CHARCOAL, align=PP_ALIGN.LEFT
        )

    # ------------------------------------------------------------------
    # SLIDE 10 — Closing / Call to Action
    # ------------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    set_bg_solid(s10, DARK_NAVY)

    add_text_box(
        s10, 'YOUR SUMMIT AWAITS',
        left=Inches(1.0), top=Inches(2.2), width=Inches(11.0), height=Inches(1.4),
        font_size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER
    )
    add_text_box(
        s10, 'Start climbing today.',
        left=Inches(3.0), top=Inches(3.8), width=Inches(7.0), height=Inches(0.8),
        font_size=28, bold=False, color=WARM_ORANGE, align=PP_ALIGN.CENTER
    )
    add_text_box(
        s10, 'Jordan Rivera  |  jordan@inspirespeaking.com  |  @jordanrises',
        left=Inches(2.5), top=Inches(6.0), width=Inches(8.0), height=Inches(0.6),
        font_size=14, bold=False, color=RGBColor(0xBB, 0xBB, 0xBB), align=PP_ALIGN.CENTER
    )

    prs.save(PPTX_PATH)
    print(f'Presentation created: {PPTX_PATH}')
    print(f'  Slides: {len(prs.slides)} (slides 1 and 6 have plain solid-color backgrounds)')


def main():
    create_video()
    create_presentation()

    # GUI-ready startup: open Inspiration_Keynote.pptx in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{PPTX_PATH}"', delay_sec=3.0)
    print('GUI_READY: LibreOffice Impress launched with Inspiration_Keynote.pptx (DISPLAY=:0)')


main()
