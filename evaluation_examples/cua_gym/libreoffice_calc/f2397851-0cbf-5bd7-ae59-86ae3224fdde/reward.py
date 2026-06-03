"""
Reward Script: Extract VLC frame at 00:30 and set as background for slides 1 and 6
Task ID: osworld_multi_apps_vlc_frame_to_slide_007
Domain: multi_apps (VLC + LibreOffice Impress)
Scoring:
  Component 1: slide_bg.png exists on the Desktop (0.30 pts)
  Component 2: Slide 1 has a full-slide background image matching slide_bg.png (0.35 pts)
  Component 3: Slide 6 has a full-slide background image matching slide_bg.png (0.35 pts)
  Total: 1.0
"""

import os
import hashlib

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user/Desktop'
PPTX_FILE = os.path.join(WORKDIR, 'Inspiration_Keynote.pptx')
SLIDE_BG_FILE = os.path.join(WORKDIR, 'slide_bg.png')

TASK_ID = 'osworld_multi_apps_vlc_frame_to_slide_007'


def verify_task():
    """
    Verify that:
    1. slide_bg.png exists on the Desktop
    2. Slides 1 and 6 in Inspiration_Keynote.pptx each have a full-slide
       picture shape whose image blob matches slide_bg.png
    Returns a progressive float between 0.0 and 1.0.
    """
    total_score = 0.0

    # Load the presentation first — if it fails, we can't verify anything
    try:
        prs = Presentation(PPTX_FILE)
    except Exception as e:
        print(f"CRITICAL: Cannot load presentation {PPTX_FILE}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: slide_bg.png exists on the Desktop (0.30 pts)
    # This is a task-introduced change: the file is extracted from the video by the agent.
    slide_bg_blob = None
    try:
        if os.path.exists(SLIDE_BG_FILE):
            with open(SLIDE_BG_FILE, 'rb') as f:
                slide_bg_blob = f.read()
            if len(slide_bg_blob) > 0:
                print(f"PASS: Component 1 — slide_bg.png exists on Desktop "
                      f"(size={len(slide_bg_blob)} bytes) (0.30 pts)")
                total_score += 0.30
            else:
                print("FAIL: Component 1 — slide_bg.png exists but is empty")
        else:
            print(f"FAIL: Component 1 — slide_bg.png not found at {SLIDE_BG_FILE}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Helper: compute MD5 of image blob
    def blob_md5(data: bytes) -> str:
        return hashlib.md5(data).hexdigest()

    # Helper: find full-slide picture shapes on a slide
    # A "background" picture covers the full slide (left=0, top=0, width=slide_width, height=slide_height)
    def find_fullslide_picture(slide, slide_width, slide_height):
        """Return the first picture shape that covers the full slide, or None."""
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Allow small tolerance (1% of dimension) for minor rounding
                w_ok = abs(shape.width - slide_width) <= slide_width * 0.01
                h_ok = abs(shape.height - slide_height) <= slide_height * 0.01
                l_ok = abs(shape.left) <= slide_width * 0.01
                t_ok = abs(shape.top) <= slide_height * 0.01
                if w_ok and h_ok and l_ok and t_ok:
                    return shape
        return None

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Component 2: Slide 1 has a full-slide background image matching slide_bg.png (0.35 pts)
    try:
        slide1 = prs.slides[0]
        pic_shape = find_fullslide_picture(slide1, slide_width, slide_height)
        if pic_shape is None:
            print("FAIL: Component 2 — Slide 1 has no full-slide picture shape")
        else:
            img_blob = pic_shape.image.blob
            if slide_bg_blob is not None:
                if blob_md5(img_blob) == blob_md5(slide_bg_blob):
                    print(f"PASS: Component 2 — Slide 1 background image matches slide_bg.png "
                          f"(md5={blob_md5(img_blob)}) (0.35 pts)")
                    total_score += 0.35
                else:
                    print(f"FAIL: Component 2 — Slide 1 image md5={blob_md5(img_blob)} "
                          f"does not match slide_bg.png md5={blob_md5(slide_bg_blob)}")
            else:
                # slide_bg.png not found; still award partial points if image covers the slide
                print(f"PASS (partial): Component 2 — Slide 1 has a full-slide picture shape "
                      f"(slide_bg.png not available for comparison) (0.35 pts)")
                total_score += 0.35
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 6 has a full-slide background image matching slide_bg.png (0.35 pts)
    try:
        slide6 = prs.slides[5]
        pic_shape = find_fullslide_picture(slide6, slide_width, slide_height)
        if pic_shape is None:
            print("FAIL: Component 3 — Slide 6 has no full-slide picture shape")
        else:
            img_blob = pic_shape.image.blob
            if slide_bg_blob is not None:
                if blob_md5(img_blob) == blob_md5(slide_bg_blob):
                    print(f"PASS: Component 3 — Slide 6 background image matches slide_bg.png "
                          f"(md5={blob_md5(img_blob)}) (0.35 pts)")
                    total_score += 0.35
                else:
                    print(f"FAIL: Component 3 — Slide 6 image md5={blob_md5(img_blob)} "
                          f"does not match slide_bg.png md5={blob_md5(slide_bg_blob)}")
            else:
                print(f"PASS (partial): Component 3 — Slide 6 has a full-slide picture shape "
                      f"(slide_bg.png not available for comparison) (0.35 pts)")
                total_score += 0.35
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


verify_task()
