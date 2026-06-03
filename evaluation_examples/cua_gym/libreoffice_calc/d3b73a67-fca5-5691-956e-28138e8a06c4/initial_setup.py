"""
Initial Setup: Desktop organizer task - messy Desktop with files from multiple projects
Task ID: osworld_multi_apps_desktop_organizer_007
Domain: os (file management)

Creates:
  - 3 empty target folders on Desktop: Project_Alpha, Project_Beta, Shared_Resources
  - 8 files scattered on Desktop (not organized yet):
      alpha_design_spec.pdf, alpha_roadmap.docx, alpha_budget.xlsx,
      beta_launch_plan.pptx, beta_user_research.xlsx, beta_wireframes.pdf,
      common_template.dotx, brand_guidelines.pdf
  - Opens Nautilus file manager on the Desktop
"""

import os
import shlex
import subprocess
import time
from pathlib import Path

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_pdf_stub(path: str, title: str, content_lines: list):
    """Create a minimal valid PDF file using basic PDF structure."""
    lines = [
        '%PDF-1.4',
        '1 0 obj',
        '<< /Type /Catalog /Pages 2 0 R >>',
        'endobj',
        '2 0 obj',
        '<< /Type /Pages /Kids [3 0 R] /Count 1 >>',
        'endobj',
        '3 0 obj',
        '<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]',
        '   /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>',
        'endobj',
    ]
    # Build stream content
    stream_parts = [f'BT /F1 12 Tf 50 750 Td ({title}) Tj']
    y = 720
    for line in content_lines:
        stream_parts.append(f'0 -20 Td ({line}) Tj')
        y -= 20
    stream_parts.append('ET')
    stream_content = '\n'.join(stream_parts)
    lines += [
        '4 0 obj',
        f'<< /Length {len(stream_content)} >>',
        'stream',
        stream_content,
        'endstream',
        'endobj',
        '5 0 obj',
        '<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>',
        'endobj',
        'xref',
        '0 6',
        '0000000000 65535 f ',
        '0000000009 00000 n ',
        '0000000068 00000 n ',
        '0000000125 00000 n ',
        '0000000274 00000 n ',
        '0000000400 00000 n ',
        'trailer',
        '<< /Size 6 /Root 1 0 R >>',
        'startxref',
        '500',
        '%%EOF',
    ]
    Path(path).write_text('\n'.join(lines))


def create_initial():
    # Ensure Desktop directory exists
    os.makedirs(DESKTOP, exist_ok=True)

    # Create the 3 project folders (empty initially)
    for folder in ['Project_Alpha', 'Project_Beta', 'Shared_Resources']:
        os.makedirs(os.path.join(DESKTOP, folder), exist_ok=True)
        print(f'Created folder: {DESKTOP}/{folder}')

    # --- Create alpha_design_spec.pdf ---
    create_pdf_stub(
        os.path.join(DESKTOP, 'alpha_design_spec.pdf'),
        'Alpha Product Design Specification',
        [
            'Version: 2.1',
            'Project: Alpha',
            'Author: Sarah Chen, Lead Designer',
            'Date: 2025-03-15',
            'Status: Draft',
            'Overview: This document outlines the design requirements',
            'for the Alpha product line including UI/UX guidelines.',
        ]
    )
    print('Created: alpha_design_spec.pdf')

    # --- Create alpha_roadmap.docx ---
    try:
        from docx import Document
        from docx.shared import Pt
        doc = Document()
        doc.add_heading('Project Alpha - Product Roadmap 2025', level=1)
        doc.add_paragraph('Owner: Marcus Johnson, Product Manager')
        doc.add_paragraph('Last Updated: 2025-02-28')
        doc.add_heading('Q1 2025 Milestones', level=2)
        doc.add_paragraph('- Complete Alpha v1.0 feature freeze by March 31')
        doc.add_paragraph('- Internal QA testing and bug triage')
        doc.add_paragraph('- Beta user testing with 50 selected enterprise clients')
        doc.add_heading('Q2 2025 Goals', level=2)
        doc.add_paragraph('- Alpha v1.0 public launch')
        doc.add_paragraph('- Onboard first 200 enterprise accounts')
        doc.add_paragraph('- Collect NPS and engagement metrics for v1.1 planning')
        doc.add_heading('Key Stakeholders', level=2)
        table = doc.add_table(rows=1, cols=3)
        table.rows[0].cells[0].text = 'Name'
        table.rows[0].cells[1].text = 'Role'
        table.rows[0].cells[2].text = 'Responsibility'
        stakeholders = [
            ('Sarah Chen', 'Design Lead', 'UI/UX and brand consistency'),
            ('Marcus Johnson', 'Product Manager', 'Roadmap and prioritization'),
            ('Kevin Park', 'Engineering Lead', 'Architecture and delivery'),
            ('Priya Sharma', 'Marketing Manager', 'Go-to-market strategy'),
        ]
        for name, role, resp in stakeholders:
            row = table.add_row()
            row.cells[0].text = name
            row.cells[1].text = role
            row.cells[2].text = resp
        doc.save(os.path.join(DESKTOP, 'alpha_roadmap.docx'))
        print('Created: alpha_roadmap.docx')
    except ImportError:
        # Fallback: write minimal OOXML stub
        Path(os.path.join(DESKTOP, 'alpha_roadmap.docx')).write_bytes(b'PK\x03\x04')
        print('Created: alpha_roadmap.docx (stub)')

    # --- Create alpha_budget.xlsx ---
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = 'Budget Overview'
        headers = ['Category', 'Q1 Budget', 'Q2 Budget', 'Q3 Budget', 'Q4 Budget', 'Annual Total']
        for col, h in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=h)
            cell.font = Font(bold=True)
            cell.fill = PatternFill('solid', fgColor='4472C4')
        budget_data = [
            ['Engineering', 125000, 130000, 128000, 135000, None],
            ['Design', 45000, 48000, 46000, 50000, None],
            ['Marketing', 62000, 75000, 80000, 70000, None],
            ['QA Testing', 28000, 30000, 32000, 29000, None],
            ['Infrastructure', 18500, 19000, 19500, 20000, None],
            ['Legal & Compliance', 12000, 12000, 15000, 12000, None],
            ['Training', 8000, 6000, 8000, 10000, None],
        ]
        for r, row_data in enumerate(budget_data, 2):
            for c, val in enumerate(row_data[:5], 1):
                ws.cell(row=r, column=c, value=val)
            # Annual total (sum of Q1-Q4)
            ws.cell(row=r, column=6, value=sum(row_data[1:5]))
        # Totals row
        ws.cell(row=9, column=1, value='TOTAL').font = Font(bold=True)
        for col in range(2, 7):
            ws.cell(row=9, column=col, value=f'=SUM({ws.cell(row=2, column=col).coordinate}:{ws.cell(row=8, column=col).coordinate})')
        ws.column_dimensions['A'].width = 22
        for col_letter in ['B', 'C', 'D', 'E', 'F']:
            ws.column_dimensions[col_letter].width = 14
        # Second sheet: Notes
        ws2 = wb.create_sheet('Notes')
        ws2['A1'] = 'Budget Notes - Project Alpha 2025'
        ws2['A2'] = 'Prepared by: Finance Team'
        ws2['A3'] = 'Approved by: Lisa Wong, CFO'
        ws2['A5'] = 'Notes:'
        ws2['A6'] = '- Engineering budget includes contractor costs for Alpha v1.0 launch'
        ws2['A7'] = '- Marketing Q2-Q3 increase supports launch campaign'
        ws2['A8'] = '- Infrastructure costs projected to stabilize in Q3'
        wb.save(os.path.join(DESKTOP, 'alpha_budget.xlsx'))
        print('Created: alpha_budget.xlsx')
    except ImportError:
        print('Warning: openpyxl not available, skipping alpha_budget.xlsx')

    # --- Create beta_launch_plan.pptx ---
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()
        # Slide 1: Title
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = 'Project Beta - Launch Plan 2025'
        slide1.placeholders[1].text = 'Prepared by: Priya Sharma, Marketing Manager\nDate: 2025-04-01'
        # Slide 2: Executive Summary
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = 'Executive Summary'
        slide2.placeholders[1].text = (
            'Target Launch Date: June 15, 2025\n'
            'Target Market: SMB segment (50-500 employees)\n'
            'Launch Regions: North America, Western Europe\n'
            'Expected First-Year Revenue: $2.4M\n'
            'Key Differentiator: AI-powered workflow automation'
        )
        # Slide 3: Go-To-Market Strategy
        slide3 = prs.slides.add_slide(prs.slide_layouts[1])
        slide3.shapes.title.text = 'Go-To-Market Strategy'
        slide3.placeholders[1].text = (
            'Phase 1 (April-May): Beta program with 100 design partners\n'
            'Phase 2 (June): Public launch with PR campaign\n'
            'Phase 3 (July-December): Sales ramp and expansion\n'
            'Channel Mix: 60% direct sales, 30% channel partners, 10% self-serve'
        )
        # Slide 4: Success Metrics
        slide4 = prs.slides.add_slide(prs.slide_layouts[1])
        slide4.shapes.title.text = 'Success Metrics'
        slide4.placeholders[1].text = (
            'Month 1: 50 paying customers, $80K ARR\n'
            'Month 3: 150 paying customers, $280K ARR\n'
            'Month 6: 400 paying customers, $750K ARR\n'
            'NPS Target: >= 45\n'
            'Churn Target: < 3% monthly'
        )
        prs.save(os.path.join(DESKTOP, 'beta_launch_plan.pptx'))
        print('Created: beta_launch_plan.pptx')
    except ImportError:
        Path(os.path.join(DESKTOP, 'beta_launch_plan.pptx')).write_bytes(b'PK\x03\x04')
        print('Created: beta_launch_plan.pptx (stub)')

    # --- Create beta_user_research.xlsx ---
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = 'Interview Summary'
        headers = ['Participant ID', 'Company Size', 'Industry', 'Role', 'Pain Points', 'Desired Features', 'Willingness to Pay', 'NPS Score']
        for col, h in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=h)
            cell.font = Font(bold=True)
            cell.fill = PatternFill('solid', fgColor='70AD47')
        research_data = [
            ['P001', '120 employees', 'Healthcare', 'Operations Manager', 'Manual data entry', 'Automated reporting', '$120/mo', 8],
            ['P002', '280 employees', 'Finance', 'IT Director', 'Siloed systems', 'API integrations', '$200/mo', 7],
            ['P003', '85 employees', 'Retail', 'CEO', 'Poor analytics', 'Real-time dashboard', '$90/mo', 9],
            ['P004', '450 employees', 'Manufacturing', 'COO', 'Compliance tracking', 'Audit trail', '$350/mo', 6],
            ['P005', '65 employees', 'Professional Services', 'Managing Partner', 'Client reporting', 'White-label reports', '$150/mo', 8],
            ['P006', '310 employees', 'Technology', 'Head of Engineering', 'Deployment bottlenecks', 'CI/CD integration', '$300/mo', 9],
            ['P007', '195 employees', 'Education', 'CTO', 'Legacy software', 'Cloud migration tools', '$180/mo', 7],
            ['P008', '420 employees', 'Logistics', 'VP Operations', 'Inventory visibility', 'Supply chain tracking', '$280/mo', 8],
            ['P009', '75 employees', 'Marketing Agency', 'Director', 'Project management', 'Task automation', '$110/mo', 9],
            ['P010', '500 employees', 'Insurance', 'CISO', 'Security compliance', 'Automated audits', '$400/mo', 7],
        ]
        for r, row_data in enumerate(research_data, 2):
            for c, val in enumerate(row_data, 1):
                ws.cell(row=r, column=c, value=val)
        ws.column_dimensions['A'].width = 14
        ws.column_dimensions['C'].width = 22
        ws.column_dimensions['E'].width = 28
        ws.column_dimensions['F'].width = 28
        # Second sheet: Quantitative Survey
        ws2 = wb.create_sheet('Survey Results')
        ws2['A1'] = 'Beta User Research - Quantitative Survey'
        ws2['A2'] = 'Sample Size: 247 respondents'
        ws2['A3'] = 'Survey Date: February 2025'
        ws2['A5'] = 'Feature Priority Rankings'
        survey_headers = ['Feature', 'Very Important', 'Somewhat Important', 'Not Important']
        for col, h in enumerate(survey_headers, 1):
            ws2.cell(row=6, column=col, value=h).font = Font(bold=True)
        survey_data = [
            ['Automated Reporting', '68%', '24%', '8%'],
            ['API Integrations', '72%', '20%', '8%'],
            ['Real-time Dashboard', '81%', '15%', '4%'],
            ['Mobile Access', '55%', '35%', '10%'],
            ['Collaboration Tools', '63%', '28%', '9%'],
            ['Export to Excel/PDF', '77%', '18%', '5%'],
        ]
        for r, row_data in enumerate(survey_data, 7):
            for c, val in enumerate(row_data, 1):
                ws2.cell(row=r, column=c, value=val)
        wb.save(os.path.join(DESKTOP, 'beta_user_research.xlsx'))
        print('Created: beta_user_research.xlsx')
    except ImportError:
        print('Warning: openpyxl not available, skipping beta_user_research.xlsx')

    # --- Create beta_wireframes.pdf ---
    create_pdf_stub(
        os.path.join(DESKTOP, 'beta_wireframes.pdf'),
        'Project Beta - UI Wireframes v1.3',
        [
            'Designer: Jordan Lee, UX Lead',
            'Date: 2025-03-20',
            'Status: Review in Progress',
            'Screens Covered:',
            '  - Dashboard Home (3 variants)',
            '  - User Onboarding Flow (5 screens)',
            '  - Settings and Preferences',
            '  - Analytics Overview Page',
            '  - Mobile Responsive Breakpoints',
            'Next Steps: Stakeholder review meeting April 5 2025',
        ]
    )
    print('Created: beta_wireframes.pdf')

    # --- Create common_template.dotx ---
    # .dotx is a Word template format; create a valid .docx-compatible stub
    try:
        from docx import Document
        doc = Document()
        doc.add_heading('Company Document Template', level=1)
        doc.add_paragraph('[Company Name] - Confidential')
        doc.add_paragraph('Document Title: [TITLE]')
        doc.add_paragraph('Version: [VERSION]')
        doc.add_paragraph('Date: [DATE]')
        doc.add_paragraph('Author: [AUTHOR]')
        doc.add_paragraph('Reviewed by: [REVIEWER]')
        doc.add_heading('1. Executive Summary', level=2)
        doc.add_paragraph('[Insert executive summary here. Limit to 2-3 sentences.]')
        doc.add_heading('2. Background', level=2)
        doc.add_paragraph('[Provide relevant background and context for the document.]')
        doc.add_heading('3. Details', level=2)
        doc.add_paragraph('[Main body content goes here.]')
        doc.add_heading('4. Next Steps', level=2)
        doc.add_paragraph('[List recommended actions, owners, and due dates.]')
        doc.add_heading('5. Appendix', level=2)
        doc.add_paragraph('[Supporting data, charts, or references.]')
        doc.save(os.path.join(DESKTOP, 'common_template.dotx'))
        print('Created: common_template.dotx')
    except ImportError:
        Path(os.path.join(DESKTOP, 'common_template.dotx')).write_bytes(b'PK\x03\x04')
        print('Created: common_template.dotx (stub)')

    # --- Create brand_guidelines.pdf ---
    create_pdf_stub(
        os.path.join(DESKTOP, 'brand_guidelines.pdf'),
        'Corporate Brand Guidelines v3.2',
        [
            'Maintained by: Creative Services Team',
            'Last Updated: 2025-01-10',
            'Applies to: All internal and external communications',
            '',
            'Logo Usage:',
            '  - Minimum size: 120px wide for digital, 1 inch for print',
            '  - Clear space: equal to height of logo mark on all sides',
            '  - Approved formats: SVG, PNG (transparent), EPS',
            '',
            'Color Palette:',
            '  - Primary: #1A3C6E (deep navy blue)',
            '  - Secondary: #F5A623 (amber)',
            '  - Neutral: #4A4A4A (charcoal)',
            '  - Background: #F8F9FA (light grey)',
            '',
            'Typography:',
            '  - Headlines: Inter Bold 28-48pt',
            '  - Body Copy: Inter Regular 10-14pt',
            '  - Captions: Inter Light 8-10pt',
        ]
    )
    print('Created: brand_guidelines.pdf')

    print('\nInitial Desktop state created successfully.')
    print('Files on Desktop (unorganized):')
    import os as _os
    for f in sorted(_os.listdir(DESKTOP)):
        fpath = _os.path.join(DESKTOP, f)
        ftype = 'DIR' if _os.path.isdir(fpath) else 'FILE'
        print(f'  [{ftype}] {f}')

    # GUI-ready: Open Nautilus file manager pointing to Desktop
    launch_gui(f'nautilus "{DESKTOP}"', delay_sec=2.0)
    print('GUI_READY: launched Nautilus showing Desktop with DISPLAY=:0')


create_initial()
