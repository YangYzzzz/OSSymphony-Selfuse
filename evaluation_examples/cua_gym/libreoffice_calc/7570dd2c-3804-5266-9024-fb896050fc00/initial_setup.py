"""
Initial Setup: Create acl2025_talk.odp (20 slides, slide 8 has placeholder) and results.xlsx
Task ID: impress_cross_acad_004
Domain: libreoffice_impress + libreoffice_calc

Creates:
  /home/user/impress_cross_acad_004_initial.odp  — presentation with placeholder
  /home/user/Documents/results.xlsx              — bar chart source data
  /home/user/Documents/acl2025_talk.odp          — symlink/copy for the agent to work on
"""

import os
import subprocess
import shutil
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_cross_acad_004'
DOCS_DIR = f'{WORKDIR}/Documents'
PPTX_TEMP = f'{WORKDIR}/{TASK_ID}_initial.pptx'
ODP_OUTPUT = f'{WORKDIR}/{TASK_ID}_initial.odp'
DOCS_ODP = f'{DOCS_DIR}/acl2025_talk.odp'
XLSX_PATH = f'{DOCS_DIR}/results.xlsx'

# Ensure Documents directory exists
os.makedirs(DOCS_DIR, exist_ok=True)


def create_presentation_pptx():
    """Create ACL 2025 talk presentation with 20 slides."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide content for 20 slides of an academic talk
    slide_contents = [
        # (title, body)
        ("Cross-Domain Knowledge Transfer for Low-Resource NLP",
         "Jane Smith, Alex Kim, Maria Rodriguez\nACL 2025, Vienna, Austria"),
        ("Outline",
         "1. Motivation\n2. Related Work\n3. Method\n4. Datasets\n5. Experimental Results\n6. Analysis\n7. Ablation Study\n8. Error Analysis\n9. Conclusion"),
        ("Motivation",
         "Low-resource languages still underrepresented\nExisting methods rely on large annotated corpora\nCross-lingual transfer offers a promising direction"),
        ("Related Work: Transfer Learning",
         "mBERT (Devlin et al., 2019)\nXLM-R (Conneau et al., 2020)\nmT5 (Xue et al., 2021)\nAdapterHub (Pfeiffer et al., 2020)"),
        ("Related Work: Cross-Domain Adaptation",
         "Domain-adversarial training (Ganin et al., 2016)\nMeta-learning approaches (Gu et al., 2018)\nPrompt-based methods (Liu et al., 2023)"),
        ("Our Method: Overview",
         "Two-stage training pipeline:\n  Stage 1: Cross-domain pre-training\n  Stage 2: Task-specific fine-tuning\nKey innovation: adaptive alignment module"),
        ("Our Method: Architecture",
         "Encoder: XLM-R large (560M params)\nAlignment Module: 4-layer transformer\nClassifier: Linear projection head\nTraining objective: Combined CE + MMD loss"),
        # Slide 8 (index 7): Experimental Results with placeholder
        ("Experimental Results",
         "RESULTS CHART HERE"),
        ("Dataset Statistics",
         "Training: 45,000 sentences across 5 domains\nDevelopment: 5,000 sentences\nTest: 10,000 sentences\nLanguages: English, German, French, Spanish, Chinese"),
        ("Baseline Systems",
         "XLM-R (zero-shot)\nMLM fine-tuned\nAdapters (Houlsby et al.)\nMAML (Finn et al., 2017)\nOur method (full)"),
        ("Results: Main Comparison",
         "Our model outperforms all baselines\nAverage improvement: +4.3 F1 points\nLargest gain on SciTech domain (+6.1)\nSmallest gain on News domain (+2.8)"),
        ("Ablation Study",
         "Removing alignment module: -3.1 F1\nRemoving Stage 1 pre-training: -2.7 F1\nUsing only 2 adapter layers: -1.2 F1\nFull model achieves best performance"),
        ("Analysis: Transfer Patterns",
         "News to SciTech: 78.3 F1\nNews to Legal: 71.2 F1\nSciTech to Medical: 80.1 F1\nLegal to Medical: 69.4 F1"),
        ("Error Analysis",
         "Most errors on: rare entities (34%)\nPronoun resolution failures (22%)\nDomain-specific terminology (28%)\nAmbiguous boundaries (16%)"),
        ("Qualitative Examples",
         "Example 1: Correct transfer to SciTech\nExample 2: Failure case in Legal domain\nExample 3: Partial success in Medical\nSee paper appendix for full examples"),
        ("Low-Resource Scenario",
         "10-shot: 61.2 F1 (baseline: 54.1)\n50-shot: 68.9 F1 (baseline: 61.3)\n100-shot: 72.4 F1 (baseline: 65.8)\nFull data: 84.7 F1 (baseline: 79.8)"),
        ("Multilingual Results",
         "English: 84.7 F1\nGerman: 81.2 F1\nFrench: 79.8 F1\nSpanish: 80.4 F1\nChinese: 76.3 F1"),
        ("Limitations and Future Work",
         "Current limitations:\n  Computational cost of Stage 1\n  Limited to token classification\nFuture directions:\n  Generation tasks\n  More languages"),
        ("Conclusion",
         "Proposed cross-domain knowledge transfer method\nAchieves state-of-the-art on 5 benchmarks\nCode and models publicly available\nGitHub: github.com/example/cross-domain-nlp"),
        ("Thank You",
         "Questions?\n\nContact: jane.smith@university.edu\nPaper: arxiv.org/abs/2025.00001\nCode: github.com/example/cross-domain-nlp"),
    ]

    for i, (title_text, body_text) in enumerate(slide_contents):
        if i == 0:
            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title_text
            try:
                slide.placeholders[1].text = body_text
            except Exception:
                pass
        else:
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title_text
            try:
                body_ph = slide.placeholders[1]
                body_ph.text = body_text
                if i == 7:
                    # Slide 8: style placeholder text visually as placeholder
                    for para in body_ph.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(24)
                            run.font.italic = True
                            run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
            except Exception:
                pass

    prs.save(PPTX_TEMP)
    print(f'PPTX created: {PPTX_TEMP}')
    return PPTX_TEMP


def convert_to_odp(src_pptx, dst_odp):
    """Convert .pptx to .odp using LibreOffice headless."""
    out_dir = os.path.dirname(dst_odp)
    result = subprocess.run(
        ['libreoffice', '--headless', '--convert-to', 'odp', '--outdir', out_dir, src_pptx],
        capture_output=True, text=True, timeout=120
    )
    print(f'LibreOffice stdout: {result.stdout.strip()}')
    if result.stderr.strip():
        print(f'LibreOffice stderr: {result.stderr.strip()}')

    # LibreOffice names output based on input basename
    base = os.path.splitext(os.path.basename(src_pptx))[0]
    auto_output = os.path.join(out_dir, base + '.odp')

    if os.path.exists(auto_output) and auto_output != dst_odp:
        shutil.move(auto_output, dst_odp)

    if os.path.exists(dst_odp):
        size = os.path.getsize(dst_odp)
        print(f'ODP created: {dst_odp} ({size} bytes)')
        return True
    else:
        print(f'ERROR: ODP not found at {dst_odp}')
        return False


def create_results_xlsx():
    """Create results.xlsx with experimental results data."""
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = 'Sheet1'

    # Headers: Dataset, Baseline, OurModel
    headers = ['Dataset', 'Baseline', 'OurModel']
    for col, h in enumerate(headers, 1):
        ws.cell(row=1, column=col, value=h)

    # 5 rows of realistic F1 scores across NLP benchmark datasets
    data = [
        ['News',     79.8, 84.7],
        ['SciTech',  78.5, 84.6],
        ['Legal',    71.2, 76.4],
        ['Medical',  75.3, 81.2],
        ['Finance',  73.6, 79.1],
    ]
    for r, row_data in enumerate(data, 2):
        for c, val in enumerate(row_data, 1):
            ws.cell(row=r, column=c, value=val)

    wb.save(XLSX_PATH)
    print(f'Results XLSX created: {XLSX_PATH}')


def main():
    # Create results.xlsx
    create_results_xlsx()

    # Create presentation as PPTX
    pptx_path = create_presentation_pptx()

    # Convert PPTX -> ODP (canonical initial file)
    success = convert_to_odp(pptx_path, ODP_OUTPUT)
    if not success:
        print('ERROR: Failed to create ODP file')
        import sys
        sys.exit(1)

    # Also place the ODP in ~/Documents/ as acl2025_talk.odp for the agent
    shutil.copy(ODP_OUTPUT, DOCS_ODP)
    print(f'Copied to: {DOCS_ODP}')

    # Clean up temp pptx
    if os.path.exists(pptx_path):
        os.remove(pptx_path)

    # Verify slide 8 has the placeholder
    print('\nVerifying slide 8 content...')
    # Re-convert to check: load the odp as pptx for inspection
    verify_pptx = f'{WORKDIR}/{TASK_ID}_verify_temp.pptx'
    result = subprocess.run(
        ['libreoffice', '--headless', '--convert-to', 'pptx',
         '--outdir', WORKDIR, ODP_OUTPUT],
        capture_output=True, text=True, timeout=60
    )
    auto_pptx = f'{WORKDIR}/{TASK_ID}_initial.pptx'
    if os.path.exists(auto_pptx):
        prs = Presentation(auto_pptx)
        slide8 = prs.slides[7]
        print(f'Slide 8 title: {slide8.shapes.title.text if slide8.shapes.title else "(no title)"}')
        for shape in slide8.shapes:
            if shape.has_text_frame:
                print(f'  Shape "{shape.name}": "{shape.text_frame.text[:60]}"')
        os.remove(auto_pptx)

    print(f'\nInitial files created:')
    print(f'  Canonical initial: {ODP_OUTPUT}')
    print(f'  Agent working file: {DOCS_ODP}')
    print(f'  Results data: {XLSX_PATH}')
    print(f'  Slide 8 placeholder: "RESULTS CHART HERE"')


main()
