"""
Initial Setup: Multi-app task - VLC playing ambient video + LibreOffice Impress with presentation
Task ID: osworld_multi_apps_misc_052
Domain: multi_apps (VLC + LibreOffice Impress)
"""

import os
import shlex
import subprocess
import time

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_misc_052'
VIDEO_PATH = f'{DESKTOP}/ambient_video.mp4'
PPTX_PATH = f'{DESKTOP}/project_presentation.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_ambient_video():
    """Create ambient_video.mp4 on the Desktop with an audio track using ffmpeg."""
    os.makedirs(DESKTOP, exist_ok=True)

    # Generate a video with a pleasant ambient sine-wave audio track
    # Video: 30-second color gradient pattern
    # Audio: 440Hz sine wave (ambient tone)
    cmd = [
        'ffmpeg', '-y',
        '-f', 'lavfi',
        '-i', 'testsrc2=duration=30:size=640x480:rate=24',
        '-f', 'lavfi',
        '-i', 'sine=frequency=440:duration=30:sample_rate=44100',
        '-c:v', 'libx264',
        '-c:a', 'aac',
        '-shortest',
        VIDEO_PATH
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        print(f'Error creating video: {result.stderr}')
        # Try alternative approach with a simple tone
        cmd2 = [
            'ffmpeg', '-y',
            '-f', 'lavfi',
            '-i', 'color=c=blue:size=640x480:duration=30:rate=24',
            '-f', 'lavfi',
            '-i', 'sine=frequency=440:duration=30',
            '-c:v', 'libx264',
            '-c:a', 'aac',
            '-shortest',
            VIDEO_PATH
        ]
        result2 = subprocess.run(cmd2, capture_output=True, text=True)
        if result2.returncode != 0:
            print(f'Error creating video (alt): {result2.stderr}')
            raise RuntimeError('Failed to create ambient video')
    print(f'Ambient video created: {VIDEO_PATH}')


def create_presentation():
    """Create project_presentation.pptx with multiple slides and no embedded audio."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q3 Business Strategy"
    slide1.placeholders[1].text = "Innovation & Growth Roadmap 2025"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Market Overview"
    p2 = tf2.add_paragraph()
    p2.text = "Product Strategy"
    p3 = tf2.add_paragraph()
    p3.text = "Financial Performance"
    p4 = tf2.add_paragraph()
    p4.text = "Team Initiatives"
    p5 = tf2.add_paragraph()
    p5.text = "Next Steps"

    # --- Slide 3: Market Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Market Overview"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Total Addressable Market: $4.2B"
    p3a = tf3.add_paragraph()
    p3a.text = "Year-over-Year Growth: 18.3%"
    p3b = tf3.add_paragraph()
    p3b.text = "Key Competitors: TechCorp, InnovateCo, FutureSys"
    p3c = tf3.add_paragraph()
    p3c.text = "Our Market Share: 12.7% (+2.1% from Q2)"

    # --- Slide 4: Product Strategy ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Product Strategy"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Phase 1: Core Platform Enhancement (Q3 2025)"
    p4a = tf4.add_paragraph()
    p4a.text = "Phase 2: Mobile Experience Revamp (Q4 2025)"
    p4b = tf4.add_paragraph()
    p4b.text = "Phase 3: AI Integration Layer (Q1 2026)"
    p4c = tf4.add_paragraph()
    p4c.text = "Investment: $2.8M across all phases"

    # --- Slide 5: Financial Performance ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Financial Performance"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Q3 Revenue: $8.4M (Target: $7.9M)"
    p5a = tf5.add_paragraph()
    p5a.text = "Gross Margin: 67.2%"
    p5b = tf5.add_paragraph()
    p5b.text = "Operating Expenses: $3.1M"
    p5c = tf5.add_paragraph()
    p5c.text = "Net Profit: $2.3M"

    prs.save(PPTX_PATH)
    print(f'Presentation created: {PPTX_PATH}')


def create_initial():
    # Step 1: Create the ambient video
    create_ambient_video()

    # Step 2: Create the presentation (no embedded audio)
    create_presentation()

    # Step 3: Kill any existing VLC or LibreOffice instances for clean start
    subprocess.run(['pkill', '-f', 'vlc'], capture_output=True)
    subprocess.run(['pkill', '-f', 'soffice'], capture_output=True)
    time.sleep(2)

    # Step 4: Launch VLC playing ambient_video.mp4
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    env["VLC_VERBOSE"] = "-1"
    subprocess.Popen(
        ['vlc', VIDEO_PATH],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(3)

    # Step 5: Launch LibreOffice Impress with project_presentation.pptx
    launch_gui(f'libreoffice --impress "{PPTX_PATH}"', delay_sec=3.0)

    print('GUI_READY: VLC playing ambient_video.mp4 and LibreOffice Impress open with project_presentation.pptx')


create_initial()
