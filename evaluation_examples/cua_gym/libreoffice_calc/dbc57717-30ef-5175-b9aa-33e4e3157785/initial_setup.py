"""
Initial Setup: Organize semester files on Desktop by course
Task ID: osworld_multi_apps_doc_desktop_organize_011
Domain: multi_apps (os + libreoffice_calc + libreoffice_writer)

Creates 24 semester files on Desktop with course code prefixes.
No course folders, no index, no summary in initial state.
"""

import os
import shlex
import subprocess
import time

WORKDIR = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_doc_desktop_organize_011'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    os.makedirs(WORKDIR, exist_ok=True)

    # Define the 24 semester files across 4 courses
    # CS101: 7 files, MATH202: 6 files, ENG301: 5 files, HIST401: 6 files

    files = [
        # CS101 - Introduction to Computer Science (7 files)
        ('CS101_Assignment1_Variables.odt', 'CS101'),
        ('CS101_Assignment2_Loops.odt', 'CS101'),
        ('CS101_Lab3_Functions.odt', 'CS101'),
        ('CS101_ProblemSet1.ods', 'CS101'),
        ('CS101_Lecture_Notes_Week5.odp', 'CS101'),
        ('CS101_Reading_AlgorithmsIntro.pdf', 'CS101'),
        ('CS101_FinalProject_Proposal.odt', 'CS101'),

        # MATH202 - Calculus II (6 files)
        ('MATH202_Homework1_Integrals.odt', 'MATH202'),
        ('MATH202_Homework2_Series.odt', 'MATH202'),
        ('MATH202_ProblemSet_Week4.ods', 'MATH202'),
        ('MATH202_MidtermReview.odp', 'MATH202'),
        ('MATH202_Reading_TaylorSeries.pdf', 'MATH202'),
        ('MATH202_FinalExam_Practice.odt', 'MATH202'),

        # ENG301 - Technical Writing (5 files)
        ('ENG301_Essay1_TechReport.odt', 'ENG301'),
        ('ENG301_Essay2_Analysis.odt', 'ENG301'),
        ('ENG301_Presentation_Research.odp', 'ENG301'),
        ('ENG301_Reading_StyleGuide.pdf', 'ENG301'),
        ('ENG301_FinalPaper_Draft.odt', 'ENG301'),

        # HIST401 - Modern History (6 files)
        ('HIST401_Essay1_WWI_Causes.odt', 'HIST401'),
        ('HIST401_Essay2_ColdWar.odt', 'HIST401'),
        ('HIST401_ResearchNotes.ods', 'HIST401'),
        ('HIST401_Presentation_Timeline.odp', 'HIST401'),
        ('HIST401_Reading_PrimarySource.pdf', 'HIST401'),
        ('HIST401_FinalThesis_Outline.odt', 'HIST401'),
    ]

    # Create each file with realistic content based on file type
    for filename, course in files:
        filepath = os.path.join(WORKDIR, filename)
        ext = os.path.splitext(filename)[1]

        if ext == '.odt':
            _create_odt_file(filepath, filename, course)
        elif ext == '.ods':
            _create_ods_file(filepath, filename, course)
        elif ext == '.odp':
            _create_odp_file(filepath, filename, course)
        elif ext == '.pdf':
            _create_pdf_file(filepath, filename, course)

        print(f'Created: {filepath}')

    print(f'\nAll 24 semester files created on Desktop.')
    print(f'Courses: CS101 (7 files), MATH202 (6 files), ENG301 (5 files), HIST401 (6 files)')

    # GUI-ready: open file manager showing Desktop
    launch_gui('nautilus "/home/user/Desktop"', delay_sec=2.0)
    print('GUI_READY: Opened Nautilus file manager showing Desktop with DISPLAY=:0')


def _create_odt_file(filepath, filename, course):
    """Create a realistic .odt file using python-docx (saved as .docx content)."""
    # We'll write a minimal ODF text file
    from docx import Document
    from docx.shared import Pt

    doc = Document()

    base = os.path.splitext(filename)[0]
    title = base.replace('_', ' ')

    doc.add_heading(title, level=1)

    course_descriptions = {
        'CS101': 'Introduction to Computer Science',
        'MATH202': 'Calculus II',
        'ENG301': 'Technical Writing',
        'HIST401': 'Modern History',
    }
    course_full = course_descriptions.get(course, course)

    doc.add_paragraph(f'Course: {course} - {course_full}')
    doc.add_paragraph(f'Semester: Spring 2025')
    doc.add_paragraph()

    if 'Assignment' in filename or 'Essay' in filename or 'Homework' in filename:
        doc.add_heading('Introduction', level=2)
        doc.add_paragraph(
            f'This document contains work submitted for {course} - {course_full}. '
            f'The assignment covers core concepts from the course curriculum.'
        )
        doc.add_heading('Main Content', level=2)
        doc.add_paragraph(
            f'The following sections outline the key points and analysis for this assignment. '
            f'References to course materials and textbook chapters are included throughout.'
        )
        doc.add_paragraph(
            f'Students are expected to apply theoretical knowledge from lectures and '
            f'practical skills developed during lab sessions.'
        )
        doc.add_heading('Conclusion', level=2)
        doc.add_paragraph(
            f'This assignment demonstrates understanding of the subject matter covered in '
            f'{course} during the Spring 2025 semester.'
        )
    elif 'Lab' in filename:
        doc.add_heading('Lab Objectives', level=2)
        doc.add_paragraph('Complete the assigned exercises and document results.')
        doc.add_heading('Procedure', level=2)
        doc.add_paragraph('Follow the step-by-step instructions provided in class.')
        doc.add_heading('Results', level=2)
        doc.add_paragraph('Results and observations from lab activities are documented below.')
    elif 'FinalProject' in filename or 'FinalPaper' in filename or 'FinalThesis' in filename or 'FinalExam' in filename:
        doc.add_heading('Overview', level=2)
        doc.add_paragraph(f'Final project/paper for {course} - {course_full}.')
        doc.add_heading('Research Questions', level=2)
        doc.add_paragraph('The following questions guide this project:')
        doc.add_paragraph('1. What are the primary themes addressed in this course?')
        doc.add_paragraph('2. How do these concepts apply in practical scenarios?')
        doc.add_heading('Methodology', level=2)
        doc.add_paragraph('Research methodology and approach are described in this section.')
    else:
        doc.add_paragraph(f'Document content for {course} - {course_full}.')

    # Save as .odt (LibreOffice compatible via python-docx)
    # Note: python-docx saves .docx format; we save with .odt extension
    # LibreOffice can open both
    doc.save(filepath)


def _create_ods_file(filepath, filename, course):
    """Create a realistic .ods file (spreadsheet) using openpyxl."""
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active

    if 'ProblemSet' in filename:
        ws.title = 'Problem Set'
        ws['A1'] = 'Problem'
        ws['B1'] = 'Description'
        ws['C1'] = 'Points'
        ws['D1'] = 'Score'
        problems = [
            (1, 'Solve the differential equation', 20, 18),
            (2, 'Apply integration by parts', 15, 14),
            (3, 'Find convergence of series', 20, 17),
            (4, 'Compute partial derivatives', 25, 22),
            (5, 'Evaluate double integral', 20, 19),
        ]
        for r, (num, desc, pts, score) in enumerate(problems, 2):
            ws.cell(row=r, column=1, value=num)
            ws.cell(row=r, column=2, value=desc)
            ws.cell(row=r, column=3, value=pts)
            ws.cell(row=r, column=4, value=score)
    elif 'ResearchNotes' in filename:
        ws.title = 'Research Notes'
        ws['A1'] = 'Topic'
        ws['B1'] = 'Source'
        ws['C1'] = 'Page'
        ws['D1'] = 'Notes'
        notes = [
            ('WWI Origins', 'Clark, Christopher - Sleepwalkers', 45, 'Alliance system failures'),
            ('Treaty of Versailles', 'MacMillan, Margaret', 120, 'Economic consequences'),
            ('Cold War Start', 'Gaddis, John Lewis', 78, 'Truman Doctrine context'),
            ('Marshall Plan', 'Hogan, Michael', 234, 'Economic recovery strategy'),
            ('Korean War', 'Cumings, Bruce', 156, 'Regional conflicts post-WWII'),
            ('Cuban Missile Crisis', 'Kennedy, Robert', 89, 'Nuclear deterrence analysis'),
        ]
        for r, row_data in enumerate(notes, 2):
            for c, val in enumerate(row_data, 1):
                ws.cell(row=r, column=c, value=val)
    else:
        ws.title = 'Data'
        ws['A1'] = 'Item'
        ws['B1'] = 'Value'
        ws['C1'] = 'Category'
        for i in range(1, 11):
            ws.cell(row=i+1, column=1, value=f'Item {i}')
            ws.cell(row=i+1, column=2, value=i * 10)
            ws.cell(row=i+1, column=3, value='Data')

    wb.save(filepath)


def _create_odp_file(filepath, filename, course):
    """Create a realistic .odp file (presentation) using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    course_descriptions = {
        'CS101': 'Introduction to Computer Science',
        'MATH202': 'Calculus II',
        'ENG301': 'Technical Writing',
        'HIST401': 'Modern History',
    }
    course_full = course_descriptions.get(course, course)
    base = os.path.splitext(filename)[0]
    title_text = base.replace('_', ' ')

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    if slide.placeholders[1]:
        slide.placeholders[1].text = f'{course} - {course_full}\nSpring 2025'

    # Content slides
    slide_layout2 = prs.slide_layouts[1]

    if 'Lecture' in filename:
        topics = [
            ('Week 5 Overview', 'Key topics covered this week\nPractical applications\nHomework review'),
            ('Main Concepts', 'Core theory and principles\nDefinitions and terminology\nExamples from textbook'),
            ('Summary', 'What we learned\nNext week preview\nAssignment due dates'),
        ]
    elif 'MidtermReview' in filename or 'Presentation' in filename:
        topics = [
            ('Review Topics', 'Chapters 1-5 covered\nKey formulas and theorems\nProblem-solving strategies'),
            ('Practice Problems', 'Sample exam questions\nStep-by-step solutions\nCommon mistakes to avoid'),
            ('Final Tips', 'Study strategies\nTime management\nResources available'),
        ]
    else:
        topics = [
            ('Introduction', f'Overview of {course_full}\nLearning objectives\nCourse structure'),
            ('Main Content', 'Key concepts and theory\nApplications and examples\nGroup discussions'),
        ]

    for title_t, content_t in topics:
        slide = prs.slides.add_slide(slide_layout2)
        slide.shapes.title.text = title_t
        slide.placeholders[1].text = content_t

    prs.save(filepath)


def _create_pdf_file(filepath, filename, course):
    """Create a realistic .pdf file using fpdf2."""
    try:
        from fpdf import FPDF

        course_descriptions = {
            'CS101': 'Introduction to Computer Science',
            'MATH202': 'Calculus II',
            'ENG301': 'Technical Writing',
            'HIST401': 'Modern History',
        }
        course_full = course_descriptions.get(course, course)
        base = os.path.splitext(filename)[0]
        title_text = base.replace('_', ' ')

        pdf = FPDF()
        pdf.add_page()
        pdf.set_font('Helvetica', 'B', 16)
        pdf.cell(0, 10, title_text[:60], ln=True, align='C')
        pdf.ln(5)
        pdf.set_font('Helvetica', '', 12)
        pdf.cell(0, 8, f'Course: {course} - {course_full}', ln=True)
        pdf.cell(0, 8, 'Semester: Spring 2025', ln=True)
        pdf.ln(5)
        pdf.set_font('Helvetica', 'B', 13)
        pdf.cell(0, 8, 'Reading Material', ln=True)
        pdf.ln(3)
        pdf.set_font('Helvetica', '', 11)

        reading_contents = {
            'CS101': [
                'Chapter 1: Introduction to Algorithms',
                'An algorithm is a step-by-step procedure for solving a problem.',
                'Key concepts: time complexity, space complexity, Big-O notation.',
                '',
                'Chapter 2: Data Structures',
                'Arrays, linked lists, stacks, queues, trees, and graphs.',
                'Understanding when to use each data structure is crucial.',
            ],
            'MATH202': [
                'Section 8.1: Introduction to Taylor Series',
                'A Taylor series is an infinite sum of terms calculated from the',
                'values of a function derivatives at a single point.',
                '',
                'Definition: f(x) = sum of f^(n)(a)/n! * (x-a)^n',
                '',
                'Common Taylor Series: e^x, sin(x), cos(x), ln(1+x)',
            ],
            'ENG301': [
                'Style Guide - Technical Writing Standards',
                '',
                '1. Clarity: Use simple, direct language.',
                '2. Conciseness: Eliminate unnecessary words.',
                '3. Active voice: Prefer active over passive constructions.',
                '4. Audience awareness: Write for your specific readers.',
                '5. Structure: Use headings, lists, and white space effectively.',
            ],
            'HIST401': [
                'Primary Source Document',
                '',
                'This reading contains excerpts from historical documents',
                'related to major events of the 20th century.',
                '',
                'Context and significance are discussed in class sessions.',
                'Students should analyze the source critically.',
            ],
        }

        lines = reading_contents.get(course, [f'Reading material for {course}'])
        for line in lines:
            if line == '':
                pdf.ln(3)
            else:
                pdf.multi_cell(0, 7, line)

        pdf.output(filepath)
    except Exception as e:
        # Fallback: create a minimal PDF manually
        print(f'Warning: fpdf2 failed for {filename}: {e}. Creating placeholder.')
        with open(filepath, 'wb') as f:
            content = f'%PDF-1.4\n% {filename}\n% Course: {course}\n'
            f.write(content.encode())


create_initial()
