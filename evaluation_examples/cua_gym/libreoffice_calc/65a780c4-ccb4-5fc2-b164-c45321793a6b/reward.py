"""
Reward Script: Format bullet points on slides 2-10 for visual hierarchy
Task ID: impress_cross_acad_054
Domain: libreoffice_impress
Scoring:
  - Component 1: First bullet on each of slides 2-10 is 18pt bold (0.5 pts)
  - Component 2: Remaining bullets on each of slides 2-10 are 16pt regular (0.5 pts)
  Total: 1.0
"""

import os
from pptx import Presentation
from pptx.util import Pt

WORKDIR = '/home/user'  # VM path — all reward scripts run on the VM
TASK_ID = 'impress_cross_acad_054'

# Slides 2-10 are index 1-9
TARGET_SLIDE_INDICES = list(range(1, 10))  # indices 1..9 (slides 2..10)


def get_content_placeholder(slide):
    """Return the content placeholder (body) shape on a slide, or None."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name != 'Title 1' and 'Title' not in shape.name:
            return shape
    # Fallback: look for any non-title placeholder with text
    for shape in slide.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            if len(tf.paragraphs) > 1:
                return shape
    return None


def get_bullet_paragraphs(shape):
    """Return a list of paragraphs that have non-empty text (bullet points)."""
    return [p for p in shape.text_frame.paragraphs if p.text.strip()]


def get_run_bold(run):
    """Normalize bold: treat None and False as 'not bold'."""
    return run.font.bold is True


def get_run_size_pt(run):
    """Return font size in points (int), or None if not set."""
    if run.font.size is None:
        return None
    return round(run.font.size / 12700)


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0

    Scoring:
      Component 1: On each of slides 2-10, the FIRST bullet paragraph has
                   bold=True (and size=18pt). 0.5 points total (partial: per slide).
      Component 2: On each of slides 2-10, ALL remaining bullet paragraphs have
                   size=16pt and bold=False. 0.5 points total (partial: per slide).
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    print(f"Loaded presentation: {num_slides} slides")

    # Verify we have at least 10 slides
    if num_slides < 10:
        print(f"FAIL: Expected at least 10 slides, got {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    # --- Component 1: First bullet is 18pt bold on slides 2-10 ---
    # (0.5 pts total across 9 slides: ~0.0556 pts per slide)
    # We check each slide and award fractional points per slide
    comp1_passes = 0
    comp1_total = len(TARGET_SLIDE_INDICES)

    print("\n--- Component 1: First bullet is 18pt bold (slides 2-10) ---")
    for slide_idx in TARGET_SLIDE_INDICES:
        slide = prs.slides[slide_idx]
        shape = get_content_placeholder(slide)
        if shape is None:
            print(f"  FAIL: Slide {slide_idx + 1} — No content placeholder found")
            continue

        bullets = get_bullet_paragraphs(shape)
        if not bullets:
            print(f"  FAIL: Slide {slide_idx + 1} — No bullet paragraphs found")
            continue

        first_bullet = bullets[0]
        runs = [r for r in first_bullet.runs if r.text.strip()]
        if not runs:
            print(f"  FAIL: Slide {slide_idx + 1} — First bullet has no runs")
            continue

        # Check all runs in the first bullet for bold and 18pt
        first_run = runs[0]
        is_bold = get_run_bold(first_run)
        size_pt = get_run_size_pt(first_run)

        if is_bold and size_pt == 18:
            print(f"  PASS: Slide {slide_idx + 1} — First bullet is bold=True, size=18pt")
            comp1_passes += 1
        elif is_bold and size_pt != 18:
            print(f"  FAIL: Slide {slide_idx + 1} — First bullet is bold=True but size={size_pt}pt (expected 18pt)")
        elif not is_bold:
            print(f"  FAIL: Slide {slide_idx + 1} — First bullet is bold=False (expected True), size={size_pt}pt")

    comp1_score = (comp1_passes / comp1_total) * 0.5
    print(f"Component 1 result: {comp1_passes}/{comp1_total} slides pass → {comp1_score:.4f} pts")
    total_score += comp1_score

    # --- Component 2: Remaining bullets are 16pt regular on slides 2-10 ---
    # (0.5 pts total across 9 slides; per-slide we check all remaining bullets)
    comp2_passes = 0
    comp2_total = len(TARGET_SLIDE_INDICES)

    print("\n--- Component 2: Remaining bullets are 16pt regular (slides 2-10) ---")
    for slide_idx in TARGET_SLIDE_INDICES:
        slide = prs.slides[slide_idx]
        shape = get_content_placeholder(slide)
        if shape is None:
            print(f"  FAIL: Slide {slide_idx + 1} — No content placeholder found")
            continue

        bullets = get_bullet_paragraphs(shape)
        if len(bullets) < 2:
            # No remaining bullets to check — partial if only 1 bullet, but task says 3-5 bullets
            print(f"  FAIL: Slide {slide_idx + 1} — Only {len(bullets)} bullet(s), expected 3-5 bullets")
            continue

        remaining_bullets = bullets[1:]  # All bullets after the first
        slide_pass = True
        fail_details = []

        for para_idx, para in enumerate(remaining_bullets, start=1):
            runs = [r for r in para.runs if r.text.strip()]
            if not runs:
                continue
            run = runs[0]
            is_bold = get_run_bold(run)
            size_pt = get_run_size_pt(run)

            if size_pt != 16:
                slide_pass = False
                fail_details.append(f"Para {para_idx + 1}: size={size_pt}pt (expected 16pt)")
            if is_bold:
                slide_pass = False
                fail_details.append(f"Para {para_idx + 1}: bold=True (expected False)")

        if slide_pass:
            print(f"  PASS: Slide {slide_idx + 1} — All {len(remaining_bullets)} remaining bullets are 16pt regular")
            comp2_passes += 1
        else:
            print(f"  FAIL: Slide {slide_idx + 1} — Issues: {'; '.join(fail_details)}")

    comp2_score = (comp2_passes / comp2_total) * 0.5
    print(f"Component 2 result: {comp2_passes}/{comp2_total} slides pass → {comp2_score:.4f} pts")
    total_score += comp2_score

    final_score = min(round(total_score, 4), 1.0)
    print(f"\nScore: {total_score:.4f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against golden file (path on VM)
file_path = f'{WORKDIR}/{TASK_ID}_initial.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
