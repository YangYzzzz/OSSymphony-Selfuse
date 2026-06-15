"""
Initial Setup: Climate change project presentation - slides 2-10 with bullets at 18pt regular
Task ID: impress_cross_acad_054
Domain: libreoffice_impress

Creates ~/Documents/climate_change_project.odp with 12 slides.
Slides 2-10 each have 3-5 bullets all at 18pt regular (NOT bold).
The task is to make first bullet 18pt bold and remaining bullets 16pt regular.
"""

import os
import subprocess
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_cross_acad_054'
PPTX_PATH = f'{WORKDIR}/{TASK_ID}_initial.pptx'
ODP_PATH = f'{WORKDIR}/Documents/climate_change_project.odp'

def set_paragraph_font(para, size_pt, bold=False):
    """Set all runs in a paragraph to the given font size and bold setting."""
    for run in para.runs:
        run.font.size = Pt(size_pt)
        run.font.bold = bold

def add_text_to_para(tf, text, size_pt, bold=False, level=0, first=False):
    """Add a paragraph with the given text and formatting."""
    if first:
        para = tf.paragraphs[0]
    else:
        para = tf.add_paragraph()
    para.text = text
    para.level = level
    for run in para.runs:
        run.font.size = Pt(size_pt)
        run.font.bold = bold
    return para


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide layouts
    title_layout = prs.slide_layouts[0]   # Title Slide
    content_layout = prs.slide_layouts[1] # Title and Content

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(title_layout)
    slide1.shapes.title.text = "Climate Change: Causes, Impacts, and Solutions"
    slide1.placeholders[1].text = "Undergraduate Research Project\nEnvironmental Science 201\nSpring 2025"

    # Slides 2-10: Content slides with 3-5 bullets at 18pt regular
    slide_data = [
        {
            "title": "The Science of Climate Change",
            "bullets": [
                "Carbon dioxide levels have risen 50% since pre-industrial times",
                "Global average temperature has increased by 1.1°C since 1880",
                "Arctic ice has declined by 13% per decade since 1979",
                "Ocean heat content has increased in all ocean basins"
            ]
        },
        {
            "title": "Greenhouse Gas Emissions",
            "bullets": [
                "Fossil fuels account for 75% of global greenhouse gas emissions",
                "Energy production is the largest single source at 34% of emissions",
                "Agriculture contributes 24% through livestock and land use change",
                "Transportation sector produces 16% of global CO2 emissions",
                "Industrial processes release methane and nitrous oxide in addition to CO2"
            ]
        },
        {
            "title": "Impacts on Weather Patterns",
            "bullets": [
                "Extreme weather events are becoming more frequent and intense",
                "Hurricane wind speeds have increased by 25% over the past 30 years",
                "Drought conditions now affect 40% more land area globally",
                "Precipitation patterns have shifted causing flooding in new regions"
            ]
        },
        {
            "title": "Ocean and Marine Ecosystems",
            "bullets": [
                "Sea levels have risen 21-24 cm since the early 20th century",
                "Ocean acidification has increased by 26% since industrialization",
                "Coral bleaching events have doubled in frequency since the 1980s",
                "Marine species are migrating poleward at 72 km per decade",
                "Fisheries productivity has declined 5% per decade in tropical regions"
            ]
        },
        {
            "title": "Impact on Biodiversity",
            "bullets": [
                "One million species face extinction due to climate-related pressures",
                "Habitat loss accelerates as temperature zones shift poleward",
                "Migratory patterns disrupted across bird and marine species",
                "Phenological mismatches reduce food availability for many animals"
            ]
        },
        {
            "title": "Human Health Consequences",
            "bullets": [
                "Heat-related mortality is projected to increase threefold by 2050",
                "Vector-borne diseases like malaria are expanding into new regions",
                "Air quality deterioration causes 7 million premature deaths annually",
                "Food insecurity affects 690 million people and will worsen with warming",
                "Mental health impacts from climate disasters cost $1 trillion annually"
            ]
        },
        {
            "title": "Renewable Energy Transitions",
            "bullets": [
                "Solar and wind power costs have fallen 90% in the last decade",
                "Renewables now provide 30% of global electricity generation",
                "Battery storage capacity must increase 20-fold by 2040 to meet demand",
                "Green hydrogen could decarbonize 20% of energy demand by 2050"
            ]
        },
        {
            "title": "Policy Frameworks and Agreements",
            "bullets": [
                "The Paris Agreement commits 196 nations to limit warming to 1.5°C",
                "Carbon pricing mechanisms now cover 23% of global emissions",
                "Nationally Determined Contributions must triple current ambition levels",
                "Climate finance flows must reach $4.35 trillion annually by 2030",
                "Loss and damage funding framework established at COP27 in 2022"
            ]
        },
        {
            "title": "Individual and Community Actions",
            "bullets": [
                "Shifting to plant-based diet reduces individual carbon footprint by 73%",
                "Electric vehicles reduce lifecycle emissions by 50-70% over gasoline cars",
                "Home energy efficiency upgrades cut household emissions by 25-30%",
                "Community solar programs enable renters to access clean energy"
            ]
        },
    ]

    for slide_info in slide_data:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = slide_info["title"]

        # Remove default content placeholder and add a text box instead
        # Use the content placeholder (index 1)
        tf = slide.placeholders[1].text_frame
        tf.word_wrap = True

        bullets = slide_info["bullets"]
        for i, bullet_text in enumerate(bullets):
            if i == 0:
                # First paragraph (already exists in placeholder)
                para = tf.paragraphs[0]
                para.text = bullet_text
                para.level = 0
            else:
                para = tf.add_paragraph()
                para.text = bullet_text
                para.level = 0
            # ALL bullets at 18pt regular (not bold) in initial state
            for run in para.runs:
                run.font.size = Pt(18)
                run.font.bold = False

    # --- Slide 11: Looking Forward ---
    slide11 = prs.slides.add_slide(content_layout)
    slide11.shapes.title.text = "Looking Forward: Climate Projections"
    tf11 = slide11.placeholders[1].text_frame
    tf11.word_wrap = True
    projections = [
        "IPCC models project 1.5-4.5°C warming by 2100 depending on emissions path",
        "Tipping points such as permafrost thaw could accelerate warming nonlinearly",
        "Net zero emissions by 2050 gives a 50% chance of limiting warming to 1.5°C",
        "Adaptation investments of $300 billion per year needed in developing nations"
    ]
    for i, text in enumerate(projections):
        if i == 0:
            para = tf11.paragraphs[0]
            para.text = text
        else:
            para = tf11.add_paragraph()
            para.text = text
        for run in para.runs:
            run.font.size = Pt(18)

    # --- Slide 12: Conclusion ---
    slide12 = prs.slides.add_slide(content_layout)
    slide12.shapes.title.text = "Conclusion and Call to Action"
    tf12 = slide12.placeholders[1].text_frame
    tf12.word_wrap = True
    conclusions = [
        "Climate change is the defining challenge of the 21st century",
        "Science provides clear guidance — action must follow understanding",
        "Every fraction of a degree matters for millions of people and species",
        "Collaboration across governments, businesses, and communities is essential"
    ]
    for i, text in enumerate(conclusions):
        if i == 0:
            para = tf12.paragraphs[0]
            para.text = text
        else:
            para = tf12.add_paragraph()
            para.text = text
        for run in para.runs:
            run.font.size = Pt(18)

    # Save as pptx first
    prs.save(PPTX_PATH)
    print(f'PPTX file created: {PPTX_PATH}')

    # Ensure Documents directory exists
    subprocess.run(['mkdir', '-p', f'{WORKDIR}/Documents'], check=True)

    # Convert pptx to odp using LibreOffice
    result = subprocess.run(
        ['libreoffice', '--headless', '--convert-to', 'odp',
         '--outdir', f'{WORKDIR}/Documents', PPTX_PATH],
        capture_output=True, text=True, timeout=60
    )
    print(f'LibreOffice conversion stdout: {result.stdout}')
    print(f'LibreOffice conversion stderr: {result.stderr}')

    # Check if conversion produced the expected file
    # LibreOffice names the output file based on the input filename
    converted_path = f'{WORKDIR}/Documents/{TASK_ID}_initial.odp'
    if os.path.exists(converted_path):
        # Rename to the expected filename
        os.rename(converted_path, ODP_PATH)
        print(f'ODP file created: {ODP_PATH}')
    else:
        # List what was created
        import glob
        files = glob.glob(f'{WORKDIR}/Documents/*.odp')
        print(f'ODP files found: {files}')
        if files:
            os.rename(files[0], ODP_PATH)
            print(f'ODP file moved to: {ODP_PATH}')
        else:
            print('ERROR: No ODP file was created by LibreOffice conversion!')
            # Fallback: copy pptx to odp path
            import shutil
            shutil.copy(PPTX_PATH, ODP_PATH)
            print(f'Fallback: copied PPTX as ODP to {ODP_PATH}')

    print(f'Initial setup complete.')
    print(f'Slides: 12 total, slides 2-10 have 3-5 bullets at 18pt regular (not bold)')


create_initial()
