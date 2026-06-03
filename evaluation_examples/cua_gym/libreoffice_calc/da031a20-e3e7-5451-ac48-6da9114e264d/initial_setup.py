"""
Initial Setup: Extract audio from VLC video and embed as LibreOffice Impress background music
Task ID: osworld_multi_apps_misc_051
Domain: multi_apps (VLC + LibreOffice Impress)
"""

import os
import shlex
import subprocess
import time

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_misc_051'
VIDEO_FILE = f'{DESKTOP}/nature_documentary.mp4'
PPTX_FILE = f'{DESKTOP}/travel_slides.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_video():
    """Create a short nature documentary-style test video with audio on the Desktop."""
    os.makedirs(DESKTOP, exist_ok=True)

    if os.path.exists(VIDEO_FILE):
        print(f'Video already exists: {VIDEO_FILE}')
        return

    # Create a test video with nature-like audio (sine wave) and color bars
    # Duration: 10 seconds, 640x480 resolution
    result = subprocess.run([
        'ffmpeg', '-y',
        '-f', 'lavfi', '-i', 'testsrc=duration=10:size=640x480:rate=24',
        '-f', 'lavfi', '-i', 'sine=frequency=440:duration=10',
        '-c:v', 'libx264', '-pix_fmt', 'yuv420p',
        '-c:a', 'aac', '-b:a', '128k',
        '-shortest',
        VIDEO_FILE
    ], capture_output=True, text=True)

    if result.returncode != 0:
        print(f'ffmpeg error: {result.stderr}')
        # Try fallback with simpler approach
        result2 = subprocess.run([
            'ffmpeg', '-y',
            '-f', 'lavfi', '-i', 'color=c=blue:size=640x480:duration=10:rate=24',
            '-f', 'lavfi', '-i', 'aevalsrc=sin(440*2*PI*t):s=44100:d=10',
            '-c:v', 'libx264', '-pix_fmt', 'yuv420p',
            '-c:a', 'aac',
            '-shortest',
            VIDEO_FILE
        ], capture_output=True, text=True)
        if result2.returncode != 0:
            print(f'ffmpeg fallback error: {result2.stderr}')
            raise RuntimeError('Could not create test video')

    print(f'Video created: {VIDEO_FILE}')


def create_presentation():
    """Create travel_slides.pptx with realistic travel content (NO audio embedded)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    if os.path.exists(PPTX_FILE):
        print(f'Presentation already exists: {PPTX_FILE}')
        return

    prs = Presentation()

    # Slide 1: Title slide - Our Journey Around the World
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Our Journey Around the World"
    slide1.placeholders[1].text = "A Photo Story of Adventure and Discovery"

    # Style title
    title_run = slide1.shapes.title.text_frame.paragraphs[0].runs[0]
    title_run.font.size = Pt(40)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0x1A, 0x5E, 0x8A)

    # Background color for slide 1
    fill1 = slide1.background.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(0xF0, 0xF8, 0xFF)

    # Slide 2: Southeast Asia
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Southeast Asia"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Vietnam - Ha Long Bay"
    p2 = tf2.add_paragraph()
    p2.text = "Thailand - Ancient Temples of Chiang Mai"
    p3 = tf2.add_paragraph()
    p3.text = "Cambodia - Angkor Wat at Sunrise"
    p4 = tf2.add_paragraph()
    p4.text = "Indonesia - Rice Terraces of Bali"
    p5 = tf2.add_paragraph()
    p5.text = "Duration: 14 days | Best Time: Nov - Feb"

    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)

    # Slide 3: South America
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "South America"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Peru - Machu Picchu Hidden Citadel"
    p3a = tf3.add_paragraph()
    p3a.text = "Brazil - Amazon Rainforest Expedition"
    p3b = tf3.add_paragraph()
    p3b.text = "Argentina - Patagonia Trekking"
    p3c = tf3.add_paragraph()
    p3c.text = "Colombia - Coffee Region Discovery"
    p3d = tf3.add_paragraph()
    p3d.text = "Duration: 21 days | Best Time: Jun - Sep"

    fill3 = slide3.background.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE1)

    # Slide 4: Europe
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "European Highlights"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Italy - Amalfi Coast & Cinque Terre"
    p4a = tf4.add_paragraph()
    p4a.text = "Greece - Santorini Island Hopping"
    p4b = tf4.add_paragraph()
    p4b.text = "Norway - Fjords and Northern Lights"
    p4c = tf4.add_paragraph()
    p4c.text = "Croatia - Dubrovnik Old Town"
    p4d = tf4.add_paragraph()
    p4d.text = "Duration: 18 days | Best Time: May - Sep"

    fill4 = slide4.background.fill
    fill4.solid()
    fill4.fore_color.rgb = RGBColor(0xF3, 0xE5, 0xF5)

    # Slide 5: Africa
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "African Safari"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Kenya - Masai Mara Great Migration"
    p5a = tf5.add_paragraph()
    p5a.text = "Tanzania - Serengeti Wildlife"
    p5b = tf5.add_paragraph()
    p5b.text = "South Africa - Cape Town & Garden Route"
    p5c = tf5.add_paragraph()
    p5c.text = "Morocco - Sahara Desert Camel Trek"
    p5d = tf5.add_paragraph()
    p5d.text = "Duration: 16 days | Best Time: Jul - Oct"

    fill5 = slide5.background.fill
    fill5.solid()
    fill5.fore_color.rgb = RGBColor(0xFB, 0xE9, 0xE7)

    prs.save(PPTX_FILE)
    print(f'Presentation created: {PPTX_FILE}')


def create_initial():
    create_video()
    create_presentation()

    # Kill any existing LibreOffice and VLC instances for clean startup
    subprocess.run(['pkill', '-f', 'soffice'], capture_output=True)
    subprocess.run(['pkill', '-f', 'vlc'], capture_output=True)
    time.sleep(2)

    # Launch VLC playing the nature documentary video
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    env["VLC_VERBOSE"] = "-1"
    subprocess.Popen(
        ['vlc', VIDEO_FILE],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(2)

    # Launch LibreOffice Impress with the travel slides
    launch_gui(f'libreoffice --impress "{PPTX_FILE}"', delay_sec=3.0)

    print('GUI_READY: VLC playing nature_documentary.mp4, LibreOffice Impress open with travel_slides.pptx')


create_initial()
