"""
Initial Setup: Gaming Club Presentation + VLC playing game trailer
Task ID: osworld_multi_apps_misc_064
Domain: multi_apps (libreoffice_impress + vlc)

Creates:
  - /home/user/Desktop/game_trailer.mp4 (synthetic video with audio track)
  - /home/user/gaming_club.pptx (gaming club presentation, no embedded audio)
  - Opens gaming_club.pptx in LibreOffice Impress
  - Launches VLC playing game_trailer.mp4
"""

import os
import shlex
import subprocess
import time

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_misc_064'
PPTX_FILE = f'{WORKDIR}/gaming_club.pptx'
VIDEO_FILE = f'{DESKTOP}/game_trailer.mp4'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_game_trailer():
    """Create a synthetic game_trailer.mp4 with audio using ffmpeg."""
    os.makedirs(DESKTOP, exist_ok=True)

    if os.path.exists(VIDEO_FILE):
        print(f'Video already exists: {VIDEO_FILE}')
        return

    # Create a 15-second video with color bars and a synthesized audio tone
    # Video: lavfi testsrc pattern
    # Audio: sine wave at 440Hz simulating game music
    result = subprocess.run([
        "ffmpeg", "-y",
        "-f", "lavfi", "-i", "testsrc=duration=15:size=640x480:rate=24",
        "-f", "lavfi", "-i", "sine=frequency=440:duration=15",
        "-c:v", "libx264", "-pix_fmt", "yuv420p",
        "-c:a", "aac", "-b:a", "128k",
        "-shortest",
        VIDEO_FILE
    ], capture_output=True, text=True)

    if result.returncode != 0:
        print(f"ffmpeg error: {result.stderr}")
        # Fallback: try with different audio source
        result2 = subprocess.run([
            "ffmpeg", "-y",
            "-f", "lavfi", "-i",
            "testsrc=duration=15:size=640x480:rate=24[v];sine=frequency=440:duration=15[a]",
            "-map", "[v]", "-map", "[a]",
            "-c:v", "libx264", "-pix_fmt", "yuv420p",
            "-c:a", "aac", "-b:a", "128k",
            VIDEO_FILE
        ], capture_output=True, text=True)
        if result2.returncode != 0:
            # Simpler fallback: video only with silent audio
            subprocess.run([
                "ffmpeg", "-y",
                "-f", "lavfi", "-i", "color=c=blue:duration=15:size=640x480:rate=24",
                "-f", "lavfi", "-i", "anullsrc=r=44100:cl=stereo",
                "-c:v", "libx264", "-pix_fmt", "yuv420p",
                "-c:a", "aac", "-b:a", "128k",
                "-t", "15",
                VIDEO_FILE
            ], check=True)

    print(f'Game trailer created: {VIDEO_FILE}')


def create_gaming_club_pptx():
    """Create a gaming club presentation without embedded audio."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    if os.path.exists(PPTX_FILE):
        print(f'PPTX already exists: {PPTX_FILE}')
        return

    prs = Presentation()

    # Slide 1: Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Gaming Club"
    slide1.placeholders[1].text = "Join the Ultimate Gaming Experience"

    # Style title
    title_shape = slide1.shapes.title
    for para in title_shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(44)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xD7, 0x00)  # Gold

    # Background for slide 1
    fill1 = slide1.background.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)  # Dark navy

    # Slide 2: About the Club
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "About Our Club"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Weekly gaming tournaments"
    p2_1 = tf2.add_paragraph()
    p2_1.text = "All genres welcome: FPS, RPG, Strategy"
    p2_2 = tf2.add_paragraph()
    p2_2.text = "Monthly LAN parties with prizes"
    p2_3 = tf2.add_paragraph()
    p2_3.text = "Access to 30+ gaming PCs"

    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0x16, 0x21, 0x3E)  # Deep blue

    # Slide 3: Upcoming Events
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Upcoming Events"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Spring Tournament 2025 — March 22"
    p3_1 = tf3.add_paragraph()
    p3_1.text = "Retro Gaming Night — April 5"
    p3_2 = tf3.add_paragraph()
    p3_2.text = "Game Jam Weekend — April 19-20"
    p3_3 = tf3.add_paragraph()
    p3_3.text = "E-Sports Qualifier — May 10"

    fill3 = slide3.background.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(0x0F, 0x3D, 0x38)  # Dark teal

    # Slide 4: How to Join
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "How to Join"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Visit us at Room 214, Student Union"
    p4_1 = tf4.add_paragraph()
    p4_1.text = "Email: gamingclub@university.edu"
    p4_2 = tf4.add_paragraph()
    p4_2.text = "Discord: discord.gg/gamingclub"
    p4_3 = tf4.add_paragraph()
    p4_3.text = "Membership fee: $10/semester"

    fill4 = slide4.background.fill
    fill4.solid()
    fill4.fore_color.rgb = RGBColor(0x2C, 0x1A, 0x47)  # Deep purple

    prs.save(PPTX_FILE)
    print(f'Gaming club presentation created: {PPTX_FILE}')


def main():
    # Step 1: Create the game trailer video
    create_game_trailer()

    # Step 2: Create the gaming club presentation
    create_gaming_club_pptx()

    # Step 3: Open gaming_club.pptx in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{PPTX_FILE}"', delay_sec=3.0)
    print('Launched LibreOffice Impress with gaming_club.pptx')

    # Step 4: Launch VLC playing game_trailer.mp4
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    env["VLC_VERBOSE"] = "-1"
    subprocess.Popen(
        ["vlc", VIDEO_FILE],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(2.0)
    print('Launched VLC playing game_trailer.mp4')

    print('GUI_READY: LibreOffice Impress and VLC launched with DISPLAY=:0')


main()
