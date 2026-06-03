"""
Initial Setup: VLC frame capture to LibreOffice Impress slide background
Task ID: osworld_multi_apps_vlc_frame_to_slide_004
Domain: multi_apps (VLC + LibreOffice Impress)
"""

import os
import shlex
import subprocess
import time

WORKDIR = '/home/user'
TASK_ID = 'osworld_multi_apps_vlc_frame_to_slide_004'
DESKTOP = f'{WORKDIR}/Desktop'
PPTX_PATH = f'{DESKTOP}/Recipe_Book.pptx'
VIDEO_PATH = f'{DESKTOP}/pasta_tutorial.mp4'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_video():
    """Create pasta_tutorial.mp4 on the Desktop using ffmpeg with cooking-like colored frames."""
    os.makedirs(DESKTOP, exist_ok=True)

    # Create a 2-minute cooking tutorial video with warm color palette
    # Uses complex filter to generate frames with varying colors simulating cooking video
    cmd = [
        'ffmpeg', '-y',
        '-f', 'lavfi',
        '-i', 'color=c=tomato:size=1280x720:rate=25:duration=120',
        '-vf', (
            "drawtext=text='Pasta Tutorial':fontcolor=white:fontsize=48:x=(w-text_w)/2:y=(h-text_h)/2,"
            "drawtext=text='Cooking Step':fontcolor=yellow:fontsize=32:x=50:y=50"
        ),
        '-c:v', 'libx264',
        '-pix_fmt', 'yuv420p',
        '-preset', 'ultrafast',
        VIDEO_PATH
    ]

    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        # Fallback: simple color video without text overlay
        cmd_simple = [
            'ffmpeg', '-y',
            '-f', 'lavfi',
            '-i', 'color=c=0x8B4513:size=1280x720:rate=25:duration=120',
            '-c:v', 'libx264',
            '-pix_fmt', 'yuv420p',
            '-preset', 'ultrafast',
            VIDEO_PATH
        ]
        subprocess.run(cmd_simple, capture_output=True, check=True)

    print(f'Video created: {VIDEO_PATH}')


def create_pptx():
    """Create Recipe_Book.pptx with 6 slides on the Desktop."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()

    # Use slide width/height from default template (10x7.5 inches)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Helper: get layout by index (0=title_slide, 1=title_content, 5=blank, 6=title_only)
    def get_layout(idx):
        return prs.slide_layouts[idx]

    # --- Slide 1: Title slide ---
    slide1 = prs.slides.add_slide(get_layout(0))
    slide1.shapes.title.text = "Italian Recipe Collection"
    try:
        slide1.placeholders[1].text = "A culinary journey through Italy"
    except (KeyError, IndexError):
        pass
    # Warm red background
    bg1 = slide1.background.fill
    bg1.solid()
    bg1.fore_color.rgb = RGBColor(0xC0, 0x39, 0x2B)

    # --- Slide 2: Bruschetta al Pomodoro ---
    slide2 = prs.slides.add_slide(get_layout(1))
    slide2.shapes.title.text = "Bruschetta al Pomodoro"
    try:
        tf2 = slide2.placeholders[1].text_frame
        tf2.text = "Ingredients:"
        from pptx.util import Pt
        p = tf2.add_paragraph()
        p.text = "- 4 slices rustic bread, toasted"
        p = tf2.add_paragraph()
        p.text = "- 3 ripe Roma tomatoes, diced"
        p = tf2.add_paragraph()
        p.text = "- 2 cloves garlic, minced"
        p = tf2.add_paragraph()
        p.text = "- Fresh basil leaves, torn"
        p = tf2.add_paragraph()
        p.text = "- 3 tbsp extra virgin olive oil"
        p = tf2.add_paragraph()
        p.text = "- Salt and pepper to taste"
    except (KeyError, IndexError):
        pass
    bg2 = slide2.background.fill
    bg2.solid()
    bg2.fore_color.rgb = RGBColor(0xF5, 0xF0, 0xE8)

    # --- Slide 3: Pasta e Fagioli ---
    slide3 = prs.slides.add_slide(get_layout(1))
    slide3.shapes.title.text = "Pasta e Fagioli"
    try:
        tf3 = slide3.placeholders[1].text_frame
        tf3.text = "A hearty Italian bean and pasta soup"
        p = tf3.add_paragraph()
        p.text = ""
        p = tf3.add_paragraph()
        p.text = "Prep time: 15 minutes"
        p = tf3.add_paragraph()
        p.text = "Cook time: 45 minutes"
        p = tf3.add_paragraph()
        p.text = "Serves: 4-6 people"
        p = tf3.add_paragraph()
        p.text = ""
        p = tf3.add_paragraph()
        p.text = "Key ingredients: cannellini beans, ditalini pasta,"
        p = tf3.add_paragraph()
        p.text = "pancetta, rosemary, vegetable broth"
    except (KeyError, IndexError):
        pass
    bg3 = slide3.background.fill
    bg3.solid()
    bg3.fore_color.rgb = RGBColor(0xF5, 0xF0, 0xE8)

    # --- Slide 4: Spaghetti Carbonara (NO background image - task target) ---
    slide4 = prs.slides.add_slide(get_layout(1))
    slide4.shapes.title.text = "Spaghetti Carbonara"
    try:
        tf4 = slide4.placeholders[1].text_frame
        tf4.text = "The authentic Roman carbonara recipe"
        p = tf4.add_paragraph()
        p.text = ""
        p = tf4.add_paragraph()
        p.text = "Ingredients (serves 2):"
        p = tf4.add_paragraph()
        p.text = "- 200g spaghetti"
        p = tf4.add_paragraph()
        p.text = "- 100g guanciale (cured pork cheek)"
        p = tf4.add_paragraph()
        p.text = "- 2 large eggs + 2 yolks"
        p = tf4.add_paragraph()
        p.text = "- 50g Pecorino Romano, grated"
        p = tf4.add_paragraph()
        p.text = "- Freshly ground black pepper"
    except (KeyError, IndexError):
        pass
    # Plain white background (no image) - this is what the agent needs to change
    bg4 = slide4.background.fill
    bg4.solid()
    bg4.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # --- Slide 5: Tiramisu ---
    slide5 = prs.slides.add_slide(get_layout(1))
    slide5.shapes.title.text = "Tiramisu"
    try:
        tf5 = slide5.placeholders[1].text_frame
        tf5.text = "Italy's beloved coffee-flavored dessert"
        p = tf5.add_paragraph()
        p.text = ""
        p = tf5.add_paragraph()
        p.text = "Ingredients:"
        p = tf5.add_paragraph()
        p.text = "- 300g ladyfinger biscuits (savoiardi)"
        p = tf5.add_paragraph()
        p.text = "- 500g mascarpone cheese"
        p = tf5.add_paragraph()
        p.text = "- 4 eggs, separated"
        p = tf5.add_paragraph()
        p.text = "- 100g sugar"
        p = tf5.add_paragraph()
        p.text = "- 250ml strong espresso, cooled"
        p = tf5.add_paragraph()
        p.text = "- 2 tbsp Marsala wine"
        p = tf5.add_paragraph()
        p.text = "- Cocoa powder for dusting"
    except (KeyError, IndexError):
        pass
    bg5 = slide5.background.fill
    bg5.solid()
    bg5.fore_color.rgb = RGBColor(0xF0, 0xE6, 0xD3)

    # --- Slide 6: Panna Cotta ---
    slide6 = prs.slides.add_slide(get_layout(1))
    slide6.shapes.title.text = "Panna Cotta"
    try:
        tf6 = slide6.placeholders[1].text_frame
        tf6.text = "A silky smooth Italian dessert"
        p = tf6.add_paragraph()
        p.text = ""
        p = tf6.add_paragraph()
        p.text = "Ingredients:"
        p = tf6.add_paragraph()
        p.text = "- 500ml heavy cream"
        p = tf6.add_paragraph()
        p.text = "- 3 tbsp sugar"
        p = tf6.add_paragraph()
        p.text = "- 1 tsp vanilla extract"
        p = tf6.add_paragraph()
        p.text = "- 2.5 tsp gelatin powder"
        p = tf6.add_paragraph()
        p.text = "Serve with: fresh berry coulis or caramel sauce"
    except (KeyError, IndexError):
        pass
    bg6 = slide6.background.fill
    bg6.solid()
    bg6.fore_color.rgb = RGBColor(0xF0, 0xE6, 0xD3)

    prs.save(PPTX_PATH)
    print(f'Presentation created: {PPTX_PATH} ({len(prs.slides)} slides)')


def main():
    # Create Desktop directory if needed
    os.makedirs(DESKTOP, exist_ok=True)

    # Create the video file
    create_video()

    # Create the presentation
    create_pptx()

    # GUI-ready startup: open Recipe_Book.pptx in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{PPTX_PATH}"', delay_sec=3.0)
    print('GUI_READY: launched LibreOffice Impress with Recipe_Book.pptx (DISPLAY=:0)')


main()
