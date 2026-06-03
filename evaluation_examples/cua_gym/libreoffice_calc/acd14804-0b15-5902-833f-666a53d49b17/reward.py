"""
Reward Script: VLC Frame to Slide Background
Task ID: osworld_multi_apps_vlc_frame_to_slide_004
Domain: multi_apps (VLC + LibreOffice Impress)
Scoring:
  Component 1 (0.4): Slide 4 in Recipe_Book.pptx has at least one PICTURE shape
  Component 2 (0.3): That picture covers (or nearly covers) the full slide dimensions
  Component 3 (0.3): VLC snapshot file exists AND the image blob in slide 4 matches the snapshot
Total: 1.0

Task: Open pasta_tutorial.mp4 in VLC, take a snapshot at 01:20, and set that snapshot
as the background image of slide 4 in Recipe_Book.pptx (saved to Desktop).
"""

import os
import hashlib

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'osworld_multi_apps_vlc_frame_to_slide_004'
PPTX_PATH = f'{WORKDIR}/Desktop/Recipe_Book.pptx'
# VLC default snapshot naming: vlc-snap-HH-MM.png or vlc-snap-HH-MM-SS.png
# For timestamp 01:20, the expected snapshot name is vlc-snap-01-20.png
VLC_SNAPSHOT_PATH = f'{WORKDIR}/vlc-snap-01-20.png'


def verify_task(pptx_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: file must exist and be loadable
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load PPTX file {pptx_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"CRITICAL: Presentation has only {len(prs.slides)} slides, expected at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide4 = prs.slides[3]  # 0-indexed, slide 4

    # Component 1: Slide 4 has at least one PICTURE shape (0.4 points)
    # In initial state: slide 4 has NO picture shapes
    # In golden state: slide 4 has one PICTURE shape
    try:
        pic_shapes = [s for s in slide4.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        if len(pic_shapes) >= 1:
            print(f"PASS: Component 1 — Slide 4 has {len(pic_shapes)} picture shape(s) (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 — Slide 4 has no picture shapes (expected at least 1)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: The picture covers the full slide (acting as background) (0.3 points)
    # The golden image is positioned at (0,0) and fills the entire slide (9144000 x 6858000 EMU)
    try:
        pic_shapes = [s for s in slide4.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        if pic_shapes:
            # Check if any picture covers at least 80% of the slide area
            slide_area = prs.slide_width * prs.slide_height
            full_slide_found = False
            for pic in pic_shapes:
                pic_area = pic.width * pic.height
                coverage = pic_area / slide_area if slide_area > 0 else 0
                # Also check it starts near the origin (top-left corner)
                at_origin = (pic.left <= prs.slide_width * 0.1 and
                             pic.top <= prs.slide_height * 0.1)
                if coverage >= 0.80 and at_origin:
                    full_slide_found = True
                    print(f"PASS: Component 2 — Picture covers {coverage*100:.1f}% of slide, "
                          f"positioned at ({pic.left},{pic.top}) (0.3 pts)")
                    break
            if full_slide_found:
                total_score += 0.3
            else:
                # Find the largest picture to report
                if pic_shapes:
                    largest = max(pic_shapes, key=lambda s: s.width * s.height)
                    area = largest.width * largest.height / slide_area
                    print(f"FAIL: Component 2 — Largest picture covers only {area*100:.1f}% of slide, "
                          f"expected >= 80% at top-left origin")
                else:
                    print("FAIL: Component 2 — No picture shapes found on slide 4")
        else:
            print("FAIL: Component 2 — No picture shapes on slide 4 (skipped due to Component 1 fail)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: VLC snapshot exists AND the image blob in slide 4 matches the snapshot (0.3 points)
    # This verifies the image actually came from a VLC snapshot of the pasta video at 01:20
    try:
        pic_shapes = [s for s in slide4.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        if not pic_shapes:
            print("FAIL: Component 3 — No picture shapes on slide 4 to compare")
        elif not os.path.exists(VLC_SNAPSHOT_PATH):
            # VLC snapshot does not exist — the image was not taken via VLC
            print(f"FAIL: Component 3 — VLC snapshot not found at {VLC_SNAPSHOT_PATH}")
        else:
            # Compare the image blob in slide 4 with the VLC snapshot file
            with open(VLC_SNAPSHOT_PATH, 'rb') as f:
                snapshot_data = f.read()
            snapshot_hash = hashlib.md5(snapshot_data).hexdigest()

            # Check if any picture blob matches the snapshot
            match_found = False
            for pic in pic_shapes:
                blob_hash = hashlib.md5(pic.image.blob).hexdigest()
                if blob_hash == snapshot_hash:
                    match_found = True
                    print(f"PASS: Component 3 — VLC snapshot exists and blob matches slide 4 image "
                          f"(MD5: {blob_hash}) (0.3 pts)")
                    break
            if match_found:
                total_score += 0.3
            else:
                print(f"FAIL: Component 3 — VLC snapshot exists but does not match any picture "
                      f"blob in slide 4 (snapshot MD5: {snapshot_hash})")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
if not os.path.exists(PPTX_PATH):
    print(f"File not found: {PPTX_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(PPTX_PATH)
