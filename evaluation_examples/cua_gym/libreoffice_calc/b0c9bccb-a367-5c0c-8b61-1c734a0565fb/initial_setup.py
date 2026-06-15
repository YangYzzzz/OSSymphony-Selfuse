"""
Initial Setup: Organize Desktop course materials into subject folders
Task ID: osworld_multi_apps_desktop_organizer_015
Domain: os

Creates:
  - 8 course material files on the Desktop
  - 3 subject folders (Computer_Science, Mathematics, Biology) - all empty
  - Opens Nautilus file manager on the Desktop
"""

import os
import shlex
import subprocess
import time
from pathlib import Path

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_desktop_organizer_015'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    # Ensure Desktop exists
    os.makedirs(DESKTOP, exist_ok=True)

    # Create subject folders (empty, as stated in task)
    for folder in ['Computer_Science', 'Mathematics', 'Biology']:
        os.makedirs(os.path.join(DESKTOP, folder), exist_ok=True)
    print('Created subject folders: Computer_Science, Mathematics, Biology')

    # --- algorithms_lecture.pdf (Computer Science) ---
    pdf_path = os.path.join(DESKTOP, 'algorithms_lecture.pdf')
    if not os.path.exists(pdf_path):
        try:
            from fpdf import FPDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font('Helvetica', 'B', 16)
            pdf.cell(0, 10, 'Algorithms - Lecture Notes', ln=True)
            pdf.set_font('Helvetica', size=12)
            pdf.cell(0, 8, 'CS 301: Data Structures and Algorithms', ln=True)
            pdf.ln(4)
            pdf.set_font('Helvetica', 'B', 13)
            pdf.cell(0, 8, 'Chapter 1: Sorting Algorithms', ln=True)
            pdf.set_font('Helvetica', size=11)
            content = [
                'Sorting is a fundamental operation in computer science.',
                'Comparison-based algorithms have a lower bound of O(n log n).',
                '',
                'QuickSort: Average O(n log n), Worst O(n^2)',
                '  - Divide and conquer strategy',
                '  - Partition around a pivot element',
                '',
                'MergeSort: Guaranteed O(n log n)',
                '  - Stable sort algorithm',
                '  - Requires O(n) extra space',
                '',
                'HeapSort: O(n log n) in-place',
                '  - Uses binary heap data structure',
                '',
                'Chapter 2: Graph Algorithms',
                'Breadth-First Search (BFS): O(V + E)',
                'Depth-First Search (DFS): O(V + E)',
                "Dijkstra's Algorithm: O((V + E) log V)",
                '',
                'Week 3 Assignment: Implement QuickSort with 3-way partitioning.',
            ]
            for line in content:
                pdf.cell(0, 6, line, ln=True)
            pdf.output(pdf_path)
        except ImportError:
            # Fallback: create a minimal valid PDF manually
            pdf_content = b"""%PDF-1.4
1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj
2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj
3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>endobj
4 0 obj<</Length 180>>stream
BT /F1 14 Tf 72 720 Td (Algorithms - Lecture Notes) Tj 0 -24 Td /F1 12 Tf (CS 301: Data Structures and Algorithms) Tj 0 -18 Td (QuickSort, MergeSort, HeapSort, Graph Algorithms) Tj ET
endstream endobj
5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj
xref
0 6
0000000000 65535 f
0000000009 00000 n
0000000058 00000 n
0000000115 00000 n
0000000266 00000 n
0000000498 00000 n
trailer<</Size 6/Root 1 0 R>>
startxref
575
%%EOF"""
            Path(pdf_path).write_bytes(pdf_content)
    print(f'Created: {pdf_path}')

    # --- linear_algebra_notes.pdf (Mathematics) ---
    pdf_path = os.path.join(DESKTOP, 'linear_algebra_notes.pdf')
    if not os.path.exists(pdf_path):
        try:
            from fpdf import FPDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font('Helvetica', 'B', 16)
            pdf.cell(0, 10, 'Linear Algebra - Course Notes', ln=True)
            pdf.set_font('Helvetica', size=12)
            pdf.cell(0, 8, 'MATH 220: Linear Algebra', ln=True)
            pdf.ln(4)
            pdf.set_font('Helvetica', 'B', 13)
            pdf.cell(0, 8, 'Chapter 3: Eigenvalues and Eigenvectors', ln=True)
            pdf.set_font('Helvetica', size=11)
            content = [
                'An eigenvector v of a matrix A satisfies Av = lambda * v.',
                'The scalar lambda is the corresponding eigenvalue.',
                '',
                'Computing Eigenvalues:',
                '  det(A - lambda*I) = 0  (characteristic equation)',
                '',
                'Properties:',
                '  - A symmetric matrix has real eigenvalues',
                '  - Eigenvectors for distinct eigenvalues are linearly independent',
                '',
                'Chapter 4: Orthogonality',
                'Vectors u and v are orthogonal if u dot v = 0.',
                'Gram-Schmidt Process: Orthogonalizes a basis.',
                '',
                'QR Decomposition: A = QR',
                '  Q: orthogonal matrix',
                '  R: upper triangular matrix',
                '',
                'Homework 5: Compute eigenvalues for 3x3 symmetric matrix.',
            ]
            for line in content:
                pdf.cell(0, 6, line, ln=True)
            pdf.output(pdf_path)
        except ImportError:
            pdf_content = b"""%PDF-1.4
1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj
2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj
3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>endobj
4 0 obj<</Length 160>>stream
BT /F1 14 Tf 72 720 Td (Linear Algebra - Course Notes) Tj 0 -24 Td /F1 12 Tf (MATH 220: Linear Algebra) Tj 0 -18 Td (Eigenvalues, Eigenvectors, Orthogonality, QR) Tj ET
endstream endobj
5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj
xref
0 6
0000000000 65535 f
0000000009 00000 n
0000000058 00000 n
0000000115 00000 n
0000000266 00000 n
0000000478 00000 n
trailer<</Size 6/Root 1 0 R>>
startxref
555
%%EOF"""
            Path(pdf_path).write_bytes(pdf_content)
    print(f'Created: {pdf_path}')

    # --- cell_biology_chapter3.pdf (Biology) ---
    pdf_path = os.path.join(DESKTOP, 'cell_biology_chapter3.pdf')
    if not os.path.exists(pdf_path):
        try:
            from fpdf import FPDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font('Helvetica', 'B', 16)
            pdf.cell(0, 10, 'Cell Biology - Chapter 3', ln=True)
            pdf.set_font('Helvetica', size=12)
            pdf.cell(0, 8, 'BIOL 201: Introduction to Cell Biology', ln=True)
            pdf.ln(4)
            pdf.set_font('Helvetica', 'B', 13)
            pdf.cell(0, 8, 'Chapter 3: Cell Membrane and Transport', ln=True)
            pdf.set_font('Helvetica', size=11)
            content = [
                'The plasma membrane is a fluid mosaic of lipids and proteins.',
                'Phospholipid bilayer: hydrophilic heads face outward, hydrophobic tails inward.',
                '',
                'Membrane Transport:',
                '  Passive transport: diffusion, osmosis, facilitated diffusion',
                '  Active transport: requires ATP, moves against concentration gradient',
                '',
                'Endocytosis and Exocytosis:',
                '  Phagocytosis: engulfing large particles',
                '  Pinocytosis: uptake of fluids',
                '  Receptor-mediated endocytosis: highly specific',
                '',
                'Channel proteins: form pores, allow ion flow (Na+, K+, Cl-)',
                'Carrier proteins: bind and change shape to transport molecules',
                '',
                'Lab Report Due: March 18 - Membrane permeability experiment.',
            ]
            for line in content:
                pdf.cell(0, 6, line, ln=True)
            pdf.output(pdf_path)
        except ImportError:
            pdf_content = b"""%PDF-1.4
1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj
2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj
3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>endobj
4 0 obj<</Length 165>>stream
BT /F1 14 Tf 72 720 Td (Cell Biology - Chapter 3) Tj 0 -24 Td /F1 12 Tf (BIOL 201: Cell Biology) Tj 0 -18 Td (Cell Membrane, Transport, Endocytosis) Tj ET
endstream endobj
5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj
xref
0 6
0000000000 65535 f
0000000009 00000 n
0000000058 00000 n
0000000115 00000 n
0000000266 00000 n
0000000483 00000 n
trailer<</Size 6/Root 1 0 R>>
startxref
560
%%EOF"""
            Path(pdf_path).write_bytes(pdf_content)
    print(f'Created: {pdf_path}')

    # --- computational_biology_intro.pdf (CS + Biology overlap) ---
    pdf_path = os.path.join(DESKTOP, 'computational_biology_intro.pdf')
    if not os.path.exists(pdf_path):
        try:
            from fpdf import FPDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font('Helvetica', 'B', 16)
            pdf.cell(0, 10, 'Introduction to Computational Biology', ln=True)
            pdf.set_font('Helvetica', size=12)
            pdf.cell(0, 8, 'CS/BIOL 350: Computational Biology', ln=True)
            pdf.ln(4)
            pdf.set_font('Helvetica', 'B', 13)
            pdf.cell(0, 8, 'Course Overview', ln=True)
            pdf.set_font('Helvetica', size=11)
            content = [
                'Computational biology applies algorithms and data science to biological problems.',
                'Topics bridge computer science and biology research.',
                '',
                'Part 1: Sequence Analysis (Computer Science)',
                '  Dynamic programming for sequence alignment (Needleman-Wunsch)',
                '  Smith-Waterman algorithm for local alignment',
                '  BLAST: heuristic search for sequence similarity',
                '',
                'Part 2: Genomics (Biology)',
                '  DNA structure: adenine, thymine, guanine, cytosine',
                '  Gene expression: transcription and translation',
                '  Genome sequencing: shotgun sequencing, assembly',
                '',
                'Part 3: Protein Structure Prediction',
                '  AlphaFold: deep learning for protein folding',
                '  Applications in drug discovery and disease research',
                '',
                'Prerequisites: Programming (Python/C++), Introductory Biology.',
            ]
            for line in content:
                pdf.cell(0, 6, line, ln=True)
            pdf.output(pdf_path)
        except ImportError:
            pdf_content = b"""%PDF-1.4
1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj
2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj
3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>endobj
4 0 obj<</Length 185>>stream
BT /F1 14 Tf 72 720 Td (Introduction to Computational Biology) Tj 0 -24 Td /F1 12 Tf (CS/BIOL 350: Computational Biology) Tj 0 -18 Td (Sequence Analysis, Genomics, Protein Structure) Tj ET
endstream endobj
5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj
xref
0 6
0000000000 65535 f
0000000009 00000 n
0000000058 00000 n
0000000115 00000 n
0000000266 00000 n
0000000503 00000 n
trailer<</Size 6/Root 1 0 R>>
startxref
580
%%EOF"""
            Path(pdf_path).write_bytes(pdf_content)
    print(f'Created: {pdf_path}')

    # --- statistics_for_cs.pdf (Mathematics + CS overlap) ---
    pdf_path = os.path.join(DESKTOP, 'statistics_for_cs.pdf')
    if not os.path.exists(pdf_path):
        try:
            from fpdf import FPDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font('Helvetica', 'B', 16)
            pdf.cell(0, 10, 'Statistics for Computer Science', ln=True)
            pdf.set_font('Helvetica', size=12)
            pdf.cell(0, 8, 'MATH/CS 280: Probability and Statistics', ln=True)
            pdf.ln(4)
            pdf.set_font('Helvetica', 'B', 13)
            pdf.cell(0, 8, 'Unit 4: Statistical Inference', ln=True)
            pdf.set_font('Helvetica', size=11)
            content = [
                'Statistical inference draws conclusions about populations from samples.',
                '',
                'Hypothesis Testing:',
                '  Null hypothesis H0 vs. alternative hypothesis H1',
                '  Type I error (false positive): reject H0 when true',
                '  Type II error (false negative): fail to reject H0 when false',
                '  p-value: probability of observing data as extreme under H0',
                '',
                'Confidence Intervals:',
                '  95% CI: range capturing true parameter 95% of the time',
                '  Formula: x_bar +/- z*(sigma/sqrt(n))',
                '',
                'Applications in Computer Science:',
                '  A/B testing for software features',
                '  Performance benchmarking and comparison',
                '  Machine learning model evaluation (precision, recall, F1)',
                '',
                'Problem Set 4 Due: Friday - Hypothesis testing exercises.',
            ]
            for line in content:
                pdf.cell(0, 6, line, ln=True)
            pdf.output(pdf_path)
        except ImportError:
            pdf_content = b"""%PDF-1.4
1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj
2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj
3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>endobj
4 0 obj<</Length 175>>stream
BT /F1 14 Tf 72 720 Td (Statistics for Computer Science) Tj 0 -24 Td /F1 12 Tf (MATH/CS 280: Probability and Statistics) Tj 0 -18 Td (Hypothesis Testing, Confidence Intervals) Tj ET
endstream endobj
5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj
xref
0 6
0000000000 65535 f
0000000009 00000 n
0000000058 00000 n
0000000115 00000 n
0000000266 00000 n
0000000493 00000 n
trailer<</Size 6/Root 1 0 R>>
startxref
570
%%EOF"""
            Path(pdf_path).write_bytes(pdf_content)
    print(f'Created: {pdf_path}')

    # --- genetics_homework.docx (Biology) ---
    docx_path = os.path.join(DESKTOP, 'genetics_homework.docx')
    if not os.path.exists(docx_path):
        try:
            from docx import Document
            from docx.shared import Pt
            doc = Document()
            doc.add_heading('Genetics Homework - Week 7', 0)
            doc.add_paragraph('BIOL 310: Genetics and Heredity')
            doc.add_paragraph('Student: [Your Name] | Due: March 20, 2025')
            doc.add_heading('Problem 1: Mendelian Inheritance', level=1)
            doc.add_paragraph(
                'In pea plants, tall (T) is dominant over short (t), and round seeds (R) '
                'are dominant over wrinkled (r). A plant with genotype TtRr is crossed with '
                'another TtRr plant.'
            )
            doc.add_paragraph('a) What fraction of offspring will be tall with round seeds?')
            doc.add_paragraph('   Answer: 9/16 (TtRr x TtRr dihybrid cross)')
            doc.add_paragraph('b) What fraction will be short with wrinkled seeds?')
            doc.add_paragraph('   Answer: 1/16 (ttRr genotype)')
            doc.add_heading('Problem 2: Sex-Linked Inheritance', level=1)
            doc.add_paragraph(
                'Color blindness is X-linked recessive. A carrier mother (X^B X^b) and '
                'a normal-vision father (X^B Y) have children.'
            )
            doc.add_paragraph('a) Probability of color-blind daughter: 0% (daughters have X^B from father)')
            doc.add_paragraph('b) Probability of color-blind son: 50%')
            doc.add_heading('Problem 3: Hardy-Weinberg Equilibrium', level=1)
            doc.add_paragraph(
                'In a population of 500 individuals, 45 are homozygous recessive (aa). '
                'Assuming Hardy-Weinberg equilibrium, calculate allele frequencies.'
            )
            doc.add_paragraph('q^2 = 45/500 = 0.09, so q = 0.3, p = 0.7')
            doc.save(docx_path)
        except ImportError:
            # Create minimal docx (zip with XML)
            import zipfile
            import io
            content_xml = '''<?xml version="1.0" encoding="UTF-8"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
<w:body>
<w:p><w:r><w:t>Genetics Homework - Week 7</w:t></w:r></w:p>
<w:p><w:r><w:t>BIOL 310: Genetics and Heredity</w:t></w:r></w:p>
<w:p><w:r><w:t>Problem 1: Mendelian Inheritance - TtRr cross</w:t></w:r></w:p>
<w:p><w:r><w:t>Problem 2: Sex-Linked Inheritance - color blindness</w:t></w:r></w:p>
<w:p><w:r><w:t>Problem 3: Hardy-Weinberg Equilibrium</w:t></w:r></w:p>
</w:body></w:document>'''
            rels_xml = '''<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>'''
            content_types = '''<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
<Default Extension="xml" ContentType="application/xml"/>
<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>'''
            buf = io.BytesIO()
            with zipfile.ZipFile(buf, 'w') as zf:
                zf.writestr('[Content_Types].xml', content_types)
                zf.writestr('_rels/.rels', rels_xml)
                zf.writestr('word/document.xml', content_xml)
            Path(docx_path).write_bytes(buf.getvalue())
    print(f'Created: {docx_path}')

    # --- data_structures.pptx (Computer Science) ---
    pptx_path = os.path.join(DESKTOP, 'data_structures.pptx')
    if not os.path.exists(pptx_path):
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            prs = Presentation()
            # Slide 1: Title
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = 'Data Structures'
            slide.placeholders[1].text = 'CS 201: Fundamental Data Structures\nProfessor Dr. Kim | Spring 2025'
            # Slide 2: Arrays and Lists
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = 'Arrays and Linked Lists'
            tf = slide.placeholders[1].text_frame
            tf.text = 'Arrays'
            tf.add_paragraph().text = '  - Fixed size, O(1) random access'
            tf.add_paragraph().text = '  - O(n) insertion/deletion'
            tf.add_paragraph().text = 'Linked Lists'
            tf.add_paragraph().text = '  - Dynamic size, O(1) insert at head'
            tf.add_paragraph().text = '  - O(n) search'
            # Slide 3: Trees
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = 'Binary Search Trees'
            tf = slide.placeholders[1].text_frame
            tf.text = 'BST Properties'
            tf.add_paragraph().text = '  - Left subtree: smaller keys'
            tf.add_paragraph().text = '  - Right subtree: larger keys'
            tf.add_paragraph().text = 'Average case: O(log n) search, insert, delete'
            tf.add_paragraph().text = 'Worst case: O(n) when unbalanced'
            # Slide 4: Hash Tables
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = 'Hash Tables'
            tf = slide.placeholders[1].text_frame
            tf.text = 'Key Features'
            tf.add_paragraph().text = '  - O(1) average case operations'
            tf.add_paragraph().text = '  - Collision resolution: chaining, open addressing'
            tf.add_paragraph().text = '  - Load factor: n/m (elements/buckets)'
            prs.save(pptx_path)
        except ImportError:
            # Create minimal pptx
            import zipfile, io
            buf = io.BytesIO()
            with zipfile.ZipFile(buf, 'w') as zf:
                zf.writestr('[Content_Types].xml', '''<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
<Default Extension="xml" ContentType="application/xml"/>
<Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
</Types>''')
                zf.writestr('_rels/.rels', '''<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
</Relationships>''')
                zf.writestr('ppt/presentation.xml', '''<?xml version="1.0" encoding="UTF-8"?>
<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
<p:sldMasterIdLst/><p:sldSz cx="9144000" cy="6858000"/><p:notesSz cx="6858000" cy="9144000"/>
</p:presentation>''')
            Path(pptx_path).write_bytes(buf.getvalue())
    print(f'Created: {pptx_path}')

    # --- calculus_problems.xlsx (Mathematics) ---
    xlsx_path = os.path.join(DESKTOP, 'calculus_problems.xlsx')
    if not os.path.exists(xlsx_path):
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = 'Problem Set'
            # Headers
            headers = ['Problem #', 'Topic', 'Description', 'Points', 'Your Answer', 'Correct Answer', 'Score']
            for col, h in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col, value=h)
                cell.font = Font(bold=True)
                cell.fill = PatternFill(start_color='4472C4', end_color='4472C4', fill_type='solid')
                cell.font = Font(bold=True, color='FFFFFF')
            # Problems
            problems = [
                (1, 'Limits', 'Find lim(x->2) of (x^2-4)/(x-2)', 5, '', 4, ''),
                (2, 'Derivatives', 'Differentiate f(x) = 3x^3 - 2x^2 + 5x - 1', 8, '', '9x^2 - 4x + 5', ''),
                (3, 'Chain Rule', 'Find d/dx of sin(x^2 + 1)', 10, '', '2x*cos(x^2 + 1)', ''),
                (4, 'Integration', 'Evaluate integral of (2x + 3) dx from 0 to 4', 8, '', 40, ''),
                (5, 'Integration By Parts', 'Evaluate integral of x*e^x dx', 12, '', 'x*e^x - e^x + C', ''),
                (6, 'Implicit Diff', 'Find dy/dx for x^2 + y^2 = 25', 10, '', '-x/y', ''),
                (7, 'Related Rates', 'Ladder sliding problem: 10ft ladder, base moves at 2ft/s', 15, '', '-3/2 ft/s at x=6', ''),
                (8, 'Optimization', 'Find maximum area rectangle with perimeter 40m', 12, '', '100 sq m (10x10)', ''),
                (9, 'Taylor Series', 'Write first 4 terms of Taylor series for e^x at x=0', 10, '', '1 + x + x^2/2 + x^3/6', ''),
                (10, 'Partial Deriv', 'Find partial derivatives of f(x,y) = x^2*y + 3xy^2', 10, '', 'fx=2xy+3y^2; fy=x^2+6xy', ''),
            ]
            for row_data in problems:
                row_num = row_data[0] + 1
                for col, val in enumerate(row_data, 1):
                    ws.cell(row=row_num, column=col, value=val)
            # Column widths
            ws.column_dimensions['A'].width = 12
            ws.column_dimensions['B'].width = 20
            ws.column_dimensions['C'].width = 45
            ws.column_dimensions['D'].width = 8
            ws.column_dimensions['E'].width = 20
            ws.column_dimensions['F'].width = 25
            ws.column_dimensions['G'].width = 10
            # Sheet 2: Formula Reference
            ws2 = wb.create_sheet('Formula Reference')
            ws2.cell(row=1, column=1, value='Formula Reference Sheet').font = Font(bold=True, size=14)
            formulas = [
                ('Power Rule', 'd/dx[x^n] = n*x^(n-1)'),
                ('Product Rule', 'd/dx[f*g] = f\'*g + f*g\''),
                ('Quotient Rule', 'd/dx[f/g] = (f\'*g - f*g\')/g^2'),
                ('Chain Rule', 'd/dx[f(g(x))] = f\'(g(x))*g\'(x)'),
                ('Fundamental Theorem', 'F\'(x) = f(x) where F(x) = integral of f(t)dt'),
            ]
            for i, (name, formula) in enumerate(formulas, 3):
                ws2.cell(row=i, column=1, value=name).font = Font(bold=True)
                ws2.cell(row=i, column=2, value=formula)
            ws2.column_dimensions['A'].width = 20
            ws2.column_dimensions['B'].width = 50
            wb.save(xlsx_path)
        except ImportError:
            # Create minimal xlsx
            import zipfile, io
            buf = io.BytesIO()
            with zipfile.ZipFile(buf, 'w') as zf:
                zf.writestr('[Content_Types].xml', '''<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
<Default Extension="xml" ContentType="application/xml"/>
<Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>
</Types>''')
            Path(xlsx_path).write_bytes(buf.getvalue())
    print(f'Created: {xlsx_path}')

    # Verify all files created
    expected_files = [
        'algorithms_lecture.pdf',
        'linear_algebra_notes.pdf',
        'cell_biology_chapter3.pdf',
        'computational_biology_intro.pdf',
        'statistics_for_cs.pdf',
        'genetics_homework.docx',
        'data_structures.pptx',
        'calculus_problems.xlsx',
        'Computer_Science',
        'Mathematics',
        'Biology',
    ]
    print('\nVerifying Desktop contents:')
    for item in expected_files:
        path = os.path.join(DESKTOP, item)
        exists = os.path.exists(path)
        print(f'  {"OK" if exists else "MISSING"}: {item}')

    # Verify subject folders are empty (no files inside yet)
    for folder in ['Computer_Science', 'Mathematics', 'Biology']:
        folder_path = os.path.join(DESKTOP, folder)
        contents = os.listdir(folder_path)
        print(f'  Folder {folder}: {len(contents)} files (should be 0)')

    # GUI-ready startup: open Nautilus file manager on Desktop
    launch_gui(f'nautilus "{DESKTOP}"', delay_sec=2.0)
    print('GUI_READY: launched Nautilus file manager on Desktop with DISPLAY=:0')


create_initial()
