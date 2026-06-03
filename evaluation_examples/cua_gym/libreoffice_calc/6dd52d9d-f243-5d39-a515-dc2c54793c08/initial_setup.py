"""
Initial Setup: Desktop cleanup task - smart organizer
Task ID: osworld_multi_apps_desktop_organizer_013
Domain: os

Creates 6 files on the Desktop with specific modification times:
  - report_final.docx         (modified today)
  - report_final_v2.docx      (same content as report_final.docx, modified today) -- duplicate
  - budget_2024.xlsx           (modified 45 days ago) -- old
  - old_notes.txt              (modified 60 days ago) -- old
  - presentation_latest.pptx  (modified 2 days ago)
  - presentation_backup.pptx  (same content as presentation_latest.pptx) -- duplicate
"""

import os
import shlex
import subprocess
import time
import struct
import datetime

DESKTOP = '/home/user/Desktop'

# -------------------------------------------------------------------------
# Helpers
# -------------------------------------------------------------------------

def set_mtime(path: str, days_ago: int):
    """Set the modification time of a file to N days ago from now."""
    target = datetime.datetime.now() - datetime.timedelta(days=days_ago)
    ts = target.timestamp()
    os.utime(path, (ts, ts))


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch a GUI app on the VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


# -------------------------------------------------------------------------
# Minimal valid DOCX bytes (minimal ZIP/OPC structure accepted by file tools)
# We use python-docx if available, otherwise create a byte-valid placeholder.
# -------------------------------------------------------------------------

def make_docx_bytes(text_content: str) -> bytes:
    """Create a minimal valid .docx file as bytes containing the given text."""
    try:
        import io
        from docx import Document
        doc = Document()
        doc.add_paragraph(text_content)
        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()
    except ImportError:
        # Fallback: write a trivially small docx-like ZIP
        import io, zipfile
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as z:
            z.writestr('[Content_Types].xml',
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
                '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
                '<Default Extension="xml" ContentType="application/xml"/>'
                '<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
                '</Types>')
            z.writestr('_rels/.rels',
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>'
                '</Relationships>')
            z.writestr('word/document.xml',
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<w:document xmlns:wpc="http://schemas.microsoft.com/office/word/2010/wordprocessingCanvas" '
                'xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
                '<w:body><w:p><w:r><w:t>' + text_content + '</w:t></w:r></w:p></w:body>'
                '</w:document>')
            z.writestr('word/_rels/document.xml.rels',
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                '</Relationships>')
        return buf.getvalue()


def make_xlsx_bytes(sheet_data: list) -> bytes:
    """Create a minimal valid .xlsx file as bytes with given rows."""
    try:
        import io
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = 'Sheet1'
        for row in sheet_data:
            ws.append(row)
        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue()
    except ImportError:
        import io, zipfile
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as z:
            z.writestr('[Content_Types].xml',
                '<?xml version="1.0" encoding="UTF-8"?>'
                '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
                '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
                '<Default Extension="xml" ContentType="application/xml"/>'
                '<Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>'
                '</Types>')
            z.writestr('_rels/.rels',
                '<?xml version="1.0" encoding="UTF-8"?>'
                '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/>'
                '</Relationships>')
            z.writestr('xl/workbook.xml',
                '<?xml version="1.0" encoding="UTF-8"?>'
                '<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
                '<sheets><sheet name="Sheet1" sheetId="1" r:id="rId1" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"/></sheets>'
                '</workbook>')
            z.writestr('xl/_rels/workbook.xml.rels',
                '<?xml version="1.0" encoding="UTF-8"?>'
                '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/>'
                '</Relationships>')
            z.writestr('xl/worksheets/sheet1.xml',
                '<?xml version="1.0" encoding="UTF-8"?>'
                '<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
                '<sheetData/></worksheet>')
        return buf.getvalue()


def make_pptx_bytes(title: str) -> bytes:
    """Create a minimal valid .pptx file as bytes with a title slide."""
    try:
        import io
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()
    except ImportError:
        import io, zipfile
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as z:
            z.writestr('[Content_Types].xml',
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
                '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
                '<Default Extension="xml" ContentType="application/xml"/>'
                '<Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>'
                '</Types>')
            z.writestr('_rels/.rels',
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>'
                '</Relationships>')
            z.writestr('ppt/presentation.xml',
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
                '</p:presentation>')
        return buf.getvalue()


# -------------------------------------------------------------------------
# Main setup
# -------------------------------------------------------------------------

def create_initial():
    os.makedirs(DESKTOP, exist_ok=True)

    # ---------- report_final.docx ----------
    # A business report — will remain on Desktop
    report_content = (
        "Q3 2024 Financial Report\n\n"
        "Executive Summary\n"
        "Total revenue for Q3 2024 reached $4.82 million, representing a 12% increase "
        "year-over-year. Operating margins improved to 18.4%, driven by efficiency gains "
        "in the logistics and procurement divisions.\n\n"
        "Key Highlights\n"
        "- Revenue: $4,820,000 (+12% YoY)\n"
        "- Operating Profit: $887,000 (+9% YoY)\n"
        "- New Clients Acquired: 34\n"
        "- Employee Headcount: 218 (up from 201)\n\n"
        "Regional Breakdown\n"
        "North America contributed 58% of total revenue ($2,795,600). APAC grew 21% to "
        "$1,108,600. EMEA remained stable at $915,800.\n\n"
        "Outlook\n"
        "Q4 pipeline looks strong with $6.1M in identified opportunities. Management "
        "expects full-year revenue to exceed $18M for the first time.\n"
    )
    docx_bytes = make_docx_bytes(report_content)

    report_path = os.path.join(DESKTOP, 'report_final.docx')
    with open(report_path, 'wb') as f:
        f.write(docx_bytes)
    set_mtime(report_path, 0)   # today

    # ---------- report_final_v2.docx ----------
    # IDENTICAL content — this is a duplicate
    report_v2_path = os.path.join(DESKTOP, 'report_final_v2.docx')
    with open(report_v2_path, 'wb') as f:
        f.write(docx_bytes)     # byte-identical
    set_mtime(report_v2_path, 0)  # today

    # ---------- budget_2024.xlsx ----------
    # Annual budget spreadsheet — old file (45 days ago)
    budget_data = [
        ['Department',    'Q1 Budget', 'Q1 Actual', 'Q2 Budget', 'Q2 Actual'],
        ['Engineering',    450000,      432000,       480000,      461000],
        ['Marketing',      180000,      195000,       200000,      212000],
        ['Sales',          220000,      238000,       230000,      245000],
        ['HR & Operations', 95000,       91000,        98000,       94000],
        ['Finance',         60000,       58000,        62000,       59000],
        ['Customer Success', 75000,       79000,        80000,       82000],
        ['R&D',            310000,      295000,       330000,      318000],
        ['Total',         1390000,     1388000,      1480000,     1471000],
    ]
    xlsx_bytes = make_xlsx_bytes(budget_data)

    budget_path = os.path.join(DESKTOP, 'budget_2024.xlsx')
    with open(budget_path, 'wb') as f:
        f.write(xlsx_bytes)
    set_mtime(budget_path, 45)   # 45 days ago

    # ---------- old_notes.txt ----------
    # Text notes — old file (60 days ago)
    old_notes_content = (
        "Meeting Notes - Project Kickoff\n"
        "Date: Approx 2 months ago\n\n"
        "Attendees: Alice Fernandez, Bob Kwan, Priya Mehta, James Okoye\n\n"
        "Agenda\n"
        "1. Project scope and deliverables\n"
        "2. Timeline and milestones\n"
        "3. Resource allocation\n"
        "4. Risks and dependencies\n\n"
        "Action Items\n"
        "- Alice: Finalize requirements doc by end of week\n"
        "- Bob: Set up CI/CD pipeline and dev environments\n"
        "- Priya: Draft initial UI mockups for client review\n"
        "- James: Confirm vendor contracts and licensing\n\n"
        "Notes\n"
        "Team agreed on a 12-week delivery timeline with bi-weekly sprint reviews. "
        "Client has requested a live demo by week 8. Budget approved for up to $320,000 "
        "with a 10% contingency.\n\n"
        "Next meeting: Scheduled for two weeks after kickoff at 10:00 AM.\n"
    )
    notes_path = os.path.join(DESKTOP, 'old_notes.txt')
    with open(notes_path, 'w') as f:
        f.write(old_notes_content)
    set_mtime(notes_path, 60)   # 60 days ago

    # ---------- presentation_latest.pptx ----------
    # Business presentation — recent (2 days ago), will remain
    pptx_bytes = make_pptx_bytes(
        "Q4 2024 Strategy Review - Leadership Summit"
    )

    pptx_latest_path = os.path.join(DESKTOP, 'presentation_latest.pptx')
    with open(pptx_latest_path, 'wb') as f:
        f.write(pptx_bytes)
    set_mtime(pptx_latest_path, 2)  # 2 days ago

    # ---------- presentation_backup.pptx ----------
    # IDENTICAL content — duplicate
    pptx_backup_path = os.path.join(DESKTOP, 'presentation_backup.pptx')
    with open(pptx_backup_path, 'wb') as f:
        f.write(pptx_bytes)      # byte-identical
    set_mtime(pptx_backup_path, 2)  # 2 days ago

    print(f'Initial Desktop files created:')
    for fname in ['report_final.docx', 'report_final_v2.docx', 'budget_2024.xlsx',
                  'old_notes.txt', 'presentation_latest.pptx', 'presentation_backup.pptx']:
        fpath = os.path.join(DESKTOP, fname)
        mtime = datetime.datetime.fromtimestamp(os.path.getmtime(fpath))
        print(f'  {fname}: {os.path.getsize(fpath)} bytes, mtime={mtime}')

    # GUI-ready startup: open Nautilus showing the Desktop
    launch_gui('nautilus "/home/user/Desktop"', delay_sec=2.0)
    print('GUI_READY: launched Nautilus showing Desktop with DISPLAY=:0')


create_initial()
