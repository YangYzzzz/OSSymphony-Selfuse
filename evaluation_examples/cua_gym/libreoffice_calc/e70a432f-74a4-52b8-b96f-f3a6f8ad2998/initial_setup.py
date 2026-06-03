"""
Initial Setup: Extract audio from VLC video and insert as background music in portfolio.pptx
Task ID: osworld_multi_apps_misc_057
Domain: multi-app (LibreOffice Impress + VLC)

Initial state:
  - portfolio.pptx open in LibreOffice Impress (8 slides, NO audio)
  - reel.mp4 on Desktop, playing in VLC
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_misc_057'
PPTX_PATH = f'{DESKTOP}/portfolio.pptx'
VIDEO_PATH = f'{DESKTOP}/reel.mp4'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env['DISPLAY'] = ':0'
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_video():
    """Create a realistic test video (reel.mp4) with audio on the Desktop."""
    os.makedirs(DESKTOP, exist_ok=True)
    if os.path.exists(VIDEO_PATH):
        print(f'Video already exists: {VIDEO_PATH}')
        return

    # Generate a video with color bars and a sine wave audio track (30 seconds)
    result = subprocess.run([
        'ffmpeg', '-y',
        '-f', 'lavfi', '-i', 'testsrc2=duration=30:size=1280x720:rate=25',
        '-f', 'lavfi', '-i', 'sine=frequency=440:duration=30:sample_rate=44100',
        '-c:v', 'libx264', '-preset', 'fast', '-crf', '23',
        '-c:a', 'aac', '-b:a', '128k',
        '-pix_fmt', 'yuv420p',
        '-shortest',
        VIDEO_PATH
    ], capture_output=True, text=True)

    if result.returncode != 0:
        print(f'ffmpeg error: {result.stderr}')
        # Fallback: try with simpler settings
        result2 = subprocess.run([
            'ffmpeg', '-y',
            '-f', 'lavfi', '-i', 'color=c=blue:duration=30:size=1280x720:rate=25',
            '-f', 'lavfi', '-i', 'anullsrc=r=44100:cl=stereo',
            '-c:v', 'libx264', '-preset', 'ultrafast', '-crf', '28',
            '-c:a', 'aac', '-b:a', '128k',
            '-t', '30',
            '-pix_fmt', 'yuv420p',
            VIDEO_PATH
        ], capture_output=True, text=True)
        if result2.returncode != 0:
            print(f'ffmpeg fallback error: {result2.stderr}')
            raise RuntimeError('Failed to create reel.mp4')
    print(f'Video created: {VIDEO_PATH}')


def create_portfolio_pptx():
    """Create a realistic 8-slide portfolio presentation (no audio)."""
    prs = Presentation()

    # Standard widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_layouts = prs.slide_layouts
    # layout 0 = Title Slide, layout 1 = Title+Content, layout 5 = Blank, layout 6 = Title Only

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(slide_layouts[0])
    slide1.shapes.title.text = 'Alexandra Rivera'
    slide1.placeholders[1].text = 'Creative Director & Visual Designer\nPortfolio 2024'
    # Background: dark navy
    fill1 = slide1.background.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)

    # --- Slide 2: About Me ---
    slide2 = prs.slides.add_slide(slide_layouts[1])
    slide2.shapes.title.text = 'About Me'
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = 'Creative director with 8+ years crafting visual identities for global brands.'
    p2 = tf2.add_paragraph()
    p2.text = 'Expertise: Brand Strategy  |  Motion Graphics  |  UI/UX Design'
    p3 = tf2.add_paragraph()
    p3.text = 'Clients: TechNova, Meridian Hotels, GreenPath Finance, Oasis Beauty'

    # --- Slide 3: Brand Identity Project ---
    slide3 = prs.slides.add_slide(slide_layouts[1])
    slide3.shapes.title.text = 'Brand Identity — TechNova Inc.'
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = 'Complete rebrand for B2B SaaS company targeting enterprise market.'
    p = tf3.add_paragraph()
    p.text = 'Deliverables: Logo system, color palette, typography guide, icon library'
    p = tf3.add_paragraph()
    p.text = 'Outcome: 42% increase in brand recognition (user survey, Q3 2023)'
    p = tf3.add_paragraph()
    p.text = 'Timeline: 12 weeks | Budget: $85,000'

    # --- Slide 4: Motion Graphics Reel ---
    slide4 = prs.slides.add_slide(slide_layouts[1])
    slide4.shapes.title.text = 'Motion Graphics — Highlights'
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = 'Award-winning motion pieces for broadcast and digital media.'
    p = tf4.add_paragraph()
    p.text = 'Meridian Hotels — Summer Campaign (2023 AICP Award finalist)'
    p = tf4.add_paragraph()
    p.text = 'GreenPath Finance — Explainer Series (3.2M YouTube views)'
    p = tf4.add_paragraph()
    p.text = 'Oasis Beauty — Product Launch Teaser (2.8M Instagram impressions)'

    # --- Slide 5: UI/UX Design ---
    slide5 = prs.slides.add_slide(slide_layouts[1])
    slide5.shapes.title.text = 'UI/UX Design — Case Study'
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = 'Redesigned Meridian Hotels mobile booking app (iOS & Android).'
    p = tf5.add_paragraph()
    p.text = 'User Research  →  Wireframing  →  Prototype  →  Testing  →  Launch'
    p = tf5.add_paragraph()
    p.text = 'Results: 38% lower booking abandonment rate, 4.7★ App Store rating'
    p = tf5.add_paragraph()
    p.text = 'Tools: Figma, Principle, Zeplin'

    # --- Slide 6: Photography & Art Direction ---
    slide6 = prs.slides.add_slide(slide_layouts[1])
    slide6.shapes.title.text = 'Photography & Art Direction'
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = 'Campaign photography and art direction for fashion and lifestyle brands.'
    p = tf6.add_paragraph()
    p.text = 'Oasis Beauty Spring 2024 Campaign — Shot in Lisbon, Portugal'
    p = tf6.add_paragraph()
    p.text = 'GreenPath Financial Wellness Series — Documentary-style portraits'
    p = tf6.add_paragraph()
    p.text = 'TechNova Developer Conference — Event and speaker photography'

    # --- Slide 7: Client Testimonials ---
    slide7 = prs.slides.add_slide(slide_layouts[1])
    slide7.shapes.title.text = 'Client Testimonials'
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = '"Alexandra\'s visual direction transformed our brand from forgettable to iconic." — Marcus Chen, CEO, TechNova Inc.'
    p = tf7.add_paragraph()
    p.text = '"Exceptional attention to detail and a true understanding of luxury hospitality aesthetics." — Sophie Laurent, VP Marketing, Meridian Hotels'
    p = tf7.add_paragraph()
    p.text = '"Delivered on time, on budget, and well above expectations." — Raj Patel, CMO, GreenPath Finance'

    # --- Slide 8: Contact & Next Steps ---
    slide8 = prs.slides.add_slide(slide_layouts[1])
    slide8.shapes.title.text = "Let's Work Together"
    tf8 = slide8.placeholders[1].text_frame
    tf8.text = 'alexandra.rivera@studio-ar.com'
    p = tf8.add_paragraph()
    p.text = 'www.studio-ar.com'
    p = tf8.add_paragraph()
    p.text = 'LinkedIn: /in/alexandra-rivera-design'
    p = tf8.add_paragraph()
    p.text = 'Available for projects starting Q2 2024'
    p = tf8.add_paragraph()
    p.text = 'Based in Barcelona, Spain — Working worldwide'

    prs.save(PPTX_PATH)
    print(f'Portfolio presentation created: {PPTX_PATH}')


def create_initial():
    os.makedirs(DESKTOP, exist_ok=True)

    # 1. Create the video file (reel.mp4) on Desktop
    create_video()

    # 2. Create portfolio.pptx on Desktop (8 slides, no audio)
    create_portfolio_pptx()

    # 3. Kill any running LibreOffice or VLC instances before launching
    subprocess.run(['pkill', '-f', 'soffice'], capture_output=True)
    subprocess.run(['pkill', '-f', 'vlc'], capture_output=True)
    time.sleep(2)

    # 4. Launch VLC playing reel.mp4
    env_vlc = os.environ.copy()
    env_vlc['DISPLAY'] = ':0'
    env_vlc['VLC_VERBOSE'] = '-1'
    subprocess.Popen(
        ['vlc', VIDEO_PATH],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env_vlc,
    )
    time.sleep(2)

    # 5. Launch LibreOffice Impress with portfolio.pptx
    launch_gui(f'libreoffice --impress "{PPTX_PATH}"', delay_sec=3.0)

    print('GUI_READY: VLC playing reel.mp4, LibreOffice Impress open with portfolio.pptx')


create_initial()
