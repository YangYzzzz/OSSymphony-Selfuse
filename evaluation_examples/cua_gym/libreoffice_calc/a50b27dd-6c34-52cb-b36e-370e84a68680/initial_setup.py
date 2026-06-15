"""
Initial Setup: Extract audio from VLC video and add as background music to presentation
Task ID: osworld_multi_apps_misc_054
Domain: multi_apps (VLC + LibreOffice Impress)

Creates:
  - /home/user/Desktop/live_concert.mp4  (test concert video with audio)
  - /home/user/Desktop/music_presentation.pptx  (slideshow WITHOUT background music)

Then launches:
  - VLC playing live_concert.mp4
  - LibreOffice Impress with music_presentation.pptx
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_misc_054'
VIDEO_FILE = f'{DESKTOP}/live_concert.mp4'
PPTX_FILE = f'{DESKTOP}/music_presentation.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_concert_video():
    """Generate a test concert video with audio using ffmpeg."""
    os.makedirs(DESKTOP, exist_ok=True)

    if os.path.exists(VIDEO_FILE):
        print(f'Video already exists: {VIDEO_FILE}')
        return

    # Generate a test video with audio: color bars + sine tone simulating concert music
    # Duration: 30 seconds
    result = subprocess.run([
        'ffmpeg', '-y',
        '-f', 'lavfi', '-i', 'testsrc=duration=30:size=640x480:rate=25',
        '-f', 'lavfi', '-i', 'sine=frequency=440:duration=30',
        '-c:v', 'libx264', '-c:a', 'aac',
        '-pix_fmt', 'yuv420p',
        '-shortest',
        VIDEO_FILE
    ], capture_output=True, text=True)

    if result.returncode != 0:
        print(f'ffmpeg error: {result.stderr}')
        # Fallback: try with different codec
        result2 = subprocess.run([
            'ffmpeg', '-y',
            '-f', 'lavfi', '-i', 'testsrc=duration=30:size=640x480:rate=25',
            '-f', 'lavfi', '-i', 'aevalsrc=sin(2*PI*440*t):s=44100:d=30',
            '-c:v', 'mpeg4', '-c:a', 'libmp3lame',
            '-pix_fmt', 'yuv420p',
            '-shortest',
            VIDEO_FILE
        ], capture_output=True, text=True)
        if result2.returncode != 0:
            print(f'ffmpeg fallback error: {result2.stderr}')
            raise RuntimeError('Could not create test video')

    print(f'Concert video created: {VIDEO_FILE}')


def create_music_presentation():
    """Create a music-themed presentation WITHOUT background audio."""
    prs = Presentation()

    # Slide 1: Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title1 = slide1.shapes.title
    title1.text = 'Live Concert Highlights 2025'
    subtitle = slide1.placeholders[1]
    subtitle.text = 'An Evening of Outstanding Live Performances'

    # Style the title
    for para in title1.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # Background color for slide 1
    fill1 = slide1.background.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(0xF5, 0xF0, 0xE8)

    # Slide 2: Setlist / Program
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    title2.text = 'Tonight\'s Setlist'

    content2 = slide2.placeholders[1]
    tf2 = content2.text_frame
    tf2.text = 'Opening Act: The Riverside Band'
    items2 = [
        'Main Performance: Alexandra & The Resonants',
        'Featured Songs:',
        '  • "Echoes of Tomorrow" (Original)',
        '  • "Rivers Deep" (Cover)',
        '  • "Midnight Serenade" (Fan Favorite)',
        '  • "Breaking Waves" (New Release)',
        'Encore: Full Orchestra Finale',
    ]
    for item in items2:
        p = tf2.add_paragraph()
        p.text = item
        p.level = 0

    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0xF0, 0xF5, 0xFF)

    # Slide 3: Venue Info
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = 'Venue & Artist Information'

    content3 = slide3.placeholders[1]
    tf3 = content3.text_frame
    tf3.text = 'Grand Concert Hall, Downtown'
    venue_info = [
        'Capacity: 2,500 seats',
        'Founded: 1987',
        '',
        'About Alexandra & The Resonants:',
        'Formed in 2018, this award-winning ensemble',
        'blends classical composition with modern sound.',
        'Albums: "First Light" (2019), "Horizon" (2022)',
        'Awards: Best Live Act 2023, 2024',
    ]
    for info in venue_info:
        p = tf3.add_paragraph()
        p.text = info

    fill3 = slide3.background.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(0xF0, 0xFF, 0xF0)

    # Slide 4: Photo Gallery placeholder
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank

    # Add a title text box
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf4 = txBox.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = 'Concert Photo Gallery'
    p4.alignment = PP_ALIGN.CENTER
    for run in p4.runs:
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # Add a placeholder box for photos
    ph_box = slide4.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4.5))
    ph_tf = ph_box.text_frame
    ph_p = ph_tf.paragraphs[0]
    ph_p.text = '[Photo Gallery - Images from the Concert]'
    ph_p.alignment = PP_ALIGN.CENTER
    for run in ph_p.runs:
        run.font.size = Pt(18)
        run.font.italic = True
        run.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)

    fill4 = slide4.background.fill
    fill4.solid()
    fill4.fore_color.rgb = RGBColor(0xEC, 0xF0, 0xF1)

    # Slide 5: Closing slide
    slide5 = prs.slides.add_slide(prs.slide_layouts[0])
    title5 = slide5.shapes.title
    title5.text = 'Thank You for Attending!'

    subtitle5 = slide5.placeholders[1]
    subtitle5.text = 'Visit our website for upcoming events\nwww.grandconcerthall.com'

    for para in title5.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(36)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    fill5 = slide5.background.fill
    fill5.solid()
    fill5.fore_color.rgb = RGBColor(0xF5, 0xF0, 0xE8)

    # IMPORTANT: Do NOT add any background audio — the task requires the agent to do this
    prs.save(PPTX_FILE)
    print(f'Presentation created (no background audio): {PPTX_FILE}')


def create_initial():
    # Step 1: Create the concert video
    create_concert_video()

    # Step 2: Create the presentation (no background audio)
    create_music_presentation()

    # Step 3: Kill any existing VLC/LibreOffice instances for clean startup
    subprocess.run(['pkill', '-f', 'vlc'], capture_output=True)
    subprocess.run(['pkill', '-f', 'soffice'], capture_output=True)
    time.sleep(2)

    # Step 4: Launch VLC playing the concert video
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    env["VLC_VERBOSE"] = "-1"
    subprocess.Popen(
        ['vlc', VIDEO_FILE],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(2)
    print('GUI_READY: VLC launched with live_concert.mp4')

    # Step 5: Launch LibreOffice Impress with the presentation
    launch_gui(f'libreoffice --impress "{PPTX_FILE}"', delay_sec=3.0)
    print('GUI_READY: LibreOffice Impress launched with music_presentation.pptx')

    print('Initial setup complete.')
    print(f'  Concert video: {VIDEO_FILE}')
    print(f'  Presentation:  {PPTX_FILE}')
    print(f'  concert_audio.wav: NOT created (task for the agent)')


create_initial()
