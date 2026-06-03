"""
Initial Setup: Desktop file organizer - client project files sorting task
Task ID: osworld_multi_apps_desktop_organizer_012
Domain: os
"""

import os
import shlex
import subprocess
import time
from pathlib import Path

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_desktop_organizer_012'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    # Ensure Desktop exists
    os.makedirs(DESKTOP, exist_ok=True)

    # Create the four client folders
    for folder in ['Client_A', 'Client_B', 'Client_C', 'Unassigned']:
        os.makedirs(os.path.join(DESKTOP, folder), exist_ok=True)
        print(f'Created folder: {folder}')

    # ---- acme_proposal.docx ----
    # A Word document for Acme Corp proposal
    acme_proposal_content = (
        b'PK\x03\x04'  # docx minimal stub; actual content written as text below
    )
    # Write as a realistic-looking minimal docx using python-docx approach
    # Since python-docx may not be available on VM, write using a text-based approach
    # We'll create simple stub binary files that look like real files to the agent
    # The important thing is they exist with realistic names and content
    _write_docx_stub(os.path.join(DESKTOP, 'acme_proposal.docx'),
                     title='Acme Corp - Project Proposal',
                     content=(
                         'Acme Corp Project Proposal\n'
                         'Date: 2025-03-15\n'
                         'Prepared by: Sarah Chen\n\n'
                         'Executive Summary\n'
                         'This proposal outlines the scope of work for Acme Corp\'s '
                         'Q2 digital transformation initiative.\n\n'
                         'Project Scope\n'
                         '- Phase 1: Infrastructure Assessment (4 weeks)\n'
                         '- Phase 2: System Integration (8 weeks)\n'
                         '- Phase 3: User Training & Rollout (4 weeks)\n\n'
                         'Budget: $125,000\n'
                         'Timeline: Q2 2025\n'
                         'Client Contact: John Martinez, Acme Corp CTO\n'
                     ))

    # ---- acme_invoice_q3.pdf ----
    _write_pdf_stub(os.path.join(DESKTOP, 'acme_invoice_q3.pdf'),
                    title='Acme Corp - Q3 Invoice',
                    content=(
                        'INVOICE\n'
                        'Invoice #: INV-2025-0342\n'
                        'Client: Acme Corp\n'
                        'Date: 2025-07-01\n'
                        'Due Date: 2025-07-31\n\n'
                        'Services Rendered:\n'
                        '  Infrastructure Assessment    $18,500\n'
                        '  System Integration Phase 1   $32,000\n'
                        '  Consulting Hours (40h)       $12,000\n'
                        'Subtotal: $62,500\n'
                        'Tax (8%): $5,000\n'
                        'Total Due: $67,500\n'
                    ))

    # ---- globex_contract.pdf ----
    _write_pdf_stub(os.path.join(DESKTOP, 'globex_contract.pdf'),
                    title='Globex Corporation - Service Contract',
                    content=(
                        'SERVICE CONTRACT\n'
                        'Contract #: GX-2025-089\n'
                        'Client: Globex Corporation\n'
                        'Effective Date: 2025-01-01\n'
                        'Term: 12 months\n\n'
                        'Scope of Services:\n'
                        'The contractor agrees to provide software development and '
                        'consulting services to Globex Corporation as detailed in '
                        'Exhibit A.\n\n'
                        'Contract Value: $285,000 annually\n'
                        'Payment Terms: Monthly invoicing, Net-30\n'
                        'Signed by: Marcus Johnson (Account Manager)\n'
                        'Client Rep: Dr. Hank Scorpio, Globex CEO\n'
                    ))

    # ---- globex_design_brief.pptx ----
    _write_pptx_stub(os.path.join(DESKTOP, 'globex_design_brief.pptx'),
                     title='Globex - Brand Redesign Brief',
                     content=(
                         'Globex Corporation Brand Redesign Brief\n'
                         'Project: Visual Identity Overhaul 2025\n'
                         'Prepared for: Globex Marketing Team\n\n'
                         'Objectives:\n'
                         '- Modernize the Globex brand identity\n'
                         '- Align visual language with new product portfolio\n'
                         '- Improve digital presence and accessibility\n\n'
                         'Deliverables:\n'
                         '- New logo suite (primary, secondary, favicon)\n'
                         '- Color palette and typography guide\n'
                         '- Website redesign mockups\n'
                         '- Social media template kit\n\n'
                         'Timeline: 10 weeks\n'
                         'Budget Allocation: $45,000\n'
                     ))

    # ---- initech_sow.docx ----
    _write_docx_stub(os.path.join(DESKTOP, 'initech_sow.docx'),
                     title='Initech - Statement of Work',
                     content=(
                         'STATEMENT OF WORK\n'
                         'SOW #: ITH-2025-014\n'
                         'Client: Initech Corporation\n'
                         'Project: Employee Portal Development\n'
                         'Date: 2025-02-10\n\n'
                         'Project Description:\n'
                         'Development of a self-service employee portal to streamline '
                         'HR processes including leave management, payroll access, and '
                         'benefits enrollment.\n\n'
                         'Technical Requirements:\n'
                         '- React.js frontend\n'
                         '- Django REST API backend\n'
                         '- PostgreSQL database\n'
                         '- SSO integration with Active Directory\n\n'
                         'Estimated Hours: 480\n'
                         'Rate: $150/hour\n'
                         'Total Estimate: $72,000\n'
                         'Project Lead: Lisa Park\n'
                     ))

    # ---- initech_meeting_notes.txt ----
    meeting_notes_path = os.path.join(DESKTOP, 'initech_meeting_notes.txt')
    Path(meeting_notes_path).write_text(
        'Meeting Notes - Initech Project Kickoff\n'
        'Date: 2025-02-15\n'
        'Attendees: Lisa Park (PM), Tom Bradley (Dev Lead), Janet Wu (Initech IT Director)\n'
        '\n'
        'Agenda Items:\n'
        '1. Project scope review\n'
        '   - Janet confirmed employee portal requirements as per SOW\n'
        '   - Prioritize leave management module for Phase 1\n'
        '   - SSO integration to be handled in Phase 2\n'
        '\n'
        '2. Timeline discussion\n'
        '   - Development start: 2025-03-01\n'
        '   - Phase 1 delivery target: 2025-05-15\n'
        '   - Full deployment: 2025-08-01\n'
        '\n'
        '3. Access requirements\n'
        '   - Initech to provision test environment credentials by 2025-02-28\n'
        '   - Need HR system API documentation\n'
        '\n'
        'Action Items:\n'
        '- Lisa: Share project timeline Gantt chart by 2025-02-20\n'
        '- Tom: Set up dev environment\n'
        '- Janet: Provide Active Directory schema docs\n'
        '\n'
        'Next meeting: 2025-03-01 10:00 AM\n',
        encoding='utf-8'
    )
    print(f'Created: initech_meeting_notes.txt')

    # ---- generic_template.dotx ----
    _write_dotx_stub(os.path.join(DESKTOP, 'generic_template.dotx'),
                     title='Generic Document Template',
                     content=(
                         'GENERAL PURPOSE DOCUMENT TEMPLATE\n'
                         'Version: 2.1\n'
                         'Last Updated: 2025-01-10\n\n'
                         '[Company Name]\n'
                         '[Document Title]\n'
                         '[Date]\n\n'
                         'Section 1: Introduction\n'
                         '[Insert introduction text here]\n\n'
                         'Section 2: Background\n'
                         '[Insert background information here]\n\n'
                         'Section 3: Details\n'
                         '[Insert detailed information here]\n\n'
                         'Section 4: Conclusion\n'
                         '[Insert conclusion here]\n\n'
                         'Prepared by: [Author Name]\n'
                         'Reviewed by: [Reviewer Name]\n'
                         'Approved by: [Approver Name]\n'
                     ))

    # ---- internal_process.pdf ----
    _write_pdf_stub(os.path.join(DESKTOP, 'internal_process.pdf'),
                    title='Internal Process Documentation',
                    content=(
                        'INTERNAL PROCESS DOCUMENT\n'
                        'Document ID: INT-PROC-2025-003\n'
                        'Classification: Internal Use Only\n'
                        'Date: 2025-01-20\n\n'
                        'Subject: Client Onboarding Process\n\n'
                        '1. Initial Contact\n'
                        '   - Receive inquiry via CRM\n'
                        '   - Assign account manager within 24 hours\n'
                        '   - Send welcome package and NDA\n\n'
                        '2. Needs Assessment\n'
                        '   - Schedule discovery call (30-60 min)\n'
                        '   - Complete intake form\n'
                        '   - Review client requirements\n\n'
                        '3. Proposal Development\n'
                        '   - Draft SOW within 5 business days\n'
                        '   - Internal review and pricing approval\n'
                        '   - Submit to client for review\n\n'
                        '4. Contract Execution\n'
                        '   - Negotiate terms as needed\n'
                        '   - Obtain legal sign-off\n'
                        '   - Execute contract and set up project\n\n'
                        'Owner: Operations Team\n'
                        'Review Cycle: Annual\n'
                    ))

    print('\nAll files created on Desktop.')
    print('Folders created: Client_A, Client_B, Client_C, Unassigned')
    print('Files on Desktop (unsorted):')
    for f in sorted(os.listdir(DESKTOP)):
        fpath = os.path.join(DESKTOP, f)
        if os.path.isfile(fpath):
            print(f'  {f}')
        else:
            print(f'  {f}/')

    # GUI-ready startup: open Nautilus file manager showing Desktop
    launch_gui('nautilus "/home/user/Desktop"', delay_sec=2.0)
    print('GUI_READY: launched Nautilus showing Desktop with DISPLAY=:0')


def _write_docx_stub(path: str, title: str, content: str):
    """Create a minimal .docx file using python-docx if available, else text stub."""
    try:
        from docx import Document
        doc = Document()
        doc.add_heading(title, level=0)
        for line in content.strip().split('\n'):
            if line.strip():
                doc.add_paragraph(line)
            else:
                doc.add_paragraph('')
        doc.save(path)
    except ImportError:
        # Fallback: write as a zip-like stub that at least has .docx extension
        # and readable content stored inside
        import zipfile
        import io
        word_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body>
    <w:p><w:r><w:t>{title}</w:t></w:r></w:p>
'''
        for line in content.strip().split('\n'):
            escaped = line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace("'", '&apos;').replace('"', '&quot;')
            word_xml += f'    <w:p><w:r><w:t>{escaped}</w:t></w:r></w:p>\n'
        word_xml += '  </w:body>\n</w:document>\n'

        rels_xml = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>'''

        content_types_xml = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>'''

        buf = io.BytesIO()
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zf:
            zf.writestr('[Content_Types].xml', content_types_xml)
            zf.writestr('_rels/.rels', rels_xml)
            zf.writestr('word/document.xml', word_xml)
        with open(path, 'wb') as f:
            f.write(buf.getvalue())
    print(f'Created: {os.path.basename(path)}')


def _write_dotx_stub(path: str, title: str, content: str):
    """Create a minimal .dotx template file."""
    try:
        from docx import Document
        from docx.opc.constants import RELATIONSHIP_TYPE as RT
        doc = Document()
        doc.add_heading(title, level=0)
        for line in content.strip().split('\n'):
            if line.strip():
                doc.add_paragraph(line)
            else:
                doc.add_paragraph('')
        # Save as docx then rename (dotx is similar format)
        import tempfile
        tmp = path + '.tmp.docx'
        doc.save(tmp)
        os.rename(tmp, path)
    except Exception:
        # Fallback: same as docx stub
        _write_docx_stub.__wrapped__ = True
        import zipfile
        import io
        word_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body>
    <w:p><w:r><w:t>{title}</w:t></w:r></w:p>
'''
        for line in content.strip().split('\n'):
            escaped = line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            word_xml += f'    <w:p><w:r><w:t>{escaped}</w:t></w:r></w:p>\n'
        word_xml += '  </w:body>\n</w:document>\n'

        content_types_xml = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.template.main+xml"/>
</Types>'''

        rels_xml = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>'''

        buf = io.BytesIO()
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zf:
            zf.writestr('[Content_Types].xml', content_types_xml)
            zf.writestr('_rels/.rels', rels_xml)
            zf.writestr('word/document.xml', word_xml)
        with open(path, 'wb') as f:
            f.write(buf.getvalue())
    print(f'Created: {os.path.basename(path)}')


def _write_pdf_stub(path: str, title: str, content: str):
    """Create a minimal PDF file using fpdf2 if available, else minimal PDF stub."""
    try:
        from fpdf import FPDF
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font('Helvetica', 'B', 14)
        pdf.multi_cell(0, 8, title)
        pdf.ln(4)
        pdf.set_font('Helvetica', '', 11)
        for line in content.strip().split('\n'):
            pdf.multi_cell(0, 6, line if line.strip() else ' ')
        pdf.output(path)
    except ImportError:
        # Minimal valid PDF structure
        pdf_content = f"""%PDF-1.4
1 0 obj
<< /Type /Catalog /Pages 2 0 R >>
endobj

2 0 obj
<< /Type /Pages /Kids [3 0 R] /Count 1 >>
endobj

3 0 obj
<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]
   /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>
endobj

4 0 obj
<< /Length 200 >>
stream
BT
/F1 14 Tf
50 750 Td
({title[:50]}) Tj
/F1 11 Tf
0 -20 Td
(Document content available upon request.) Tj
ET
endstream
endobj

5 0 obj
<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>
endobj

xref
0 6
0000000000 65535 f
trailer
<< /Size 6 /Root 1 0 R >>
startxref
0
%%EOF
"""
        with open(path, 'w') as f:
            f.write(pdf_content)
    print(f'Created: {os.path.basename(path)}')


def _write_pptx_stub(path: str, title: str, content: str):
    """Create a minimal .pptx file using python-pptx if available."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()
        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if slide.placeholders[1]:
            slide.placeholders[1].text = 'Prepared by: Design Team\nDate: 2025-03-01'
        # Content slide
        slide_layout2 = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(slide_layout2)
        slide2.shapes.title.text = 'Project Overview'
        tf = slide2.placeholders[1].text_frame
        tf.text = content[:500]
        prs.save(path)
    except ImportError:
        import zipfile
        import io
        # Minimal pptx structure
        content_types_xml = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
</Types>'''
        ppt_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:sldMasterIdLst/>
  <p:sldSz cx="9144000" cy="6858000"/>
  <p:notesSz cx="6858000" cy="9144000"/>
</p:presentation>'''
        rels_xml = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
</Relationships>'''
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zf:
            zf.writestr('[Content_Types].xml', content_types_xml)
            zf.writestr('_rels/.rels', rels_xml)
            zf.writestr('ppt/presentation.xml', ppt_xml)
        with open(path, 'wb') as f:
            f.write(buf.getvalue())
    print(f'Created: {os.path.basename(path)}')


create_initial()
