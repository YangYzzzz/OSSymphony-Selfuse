"""
Initial Setup: Create technical_presentation.pptx with 10 slides
Task ID: osworld_multi_apps_impress_text_to_writer_008
Domain: libreoffice_impress + libreoffice_writer
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_multi_apps_impress_text_to_writer_008'
DOCS_DIR = f'{WORKDIR}/Documents'
PPTX_OUTPUT = f'{DOCS_DIR}/technical_presentation.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, text, left, top, width, height,
                 font_name='Calibri', font_size=14, bold=False, color=None):
    """Add a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullet_slide(prs, title_text, bullets):
    """Add a slide with title and multi-level bullet points."""
    slide_layout = prs.slide_layouts[1]  # Title+Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text

    tf = slide.placeholders[1].text_frame
    tf.clear()

    for i, (level, text) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = level

    return slide


def create_initial():
    os.makedirs(DOCS_DIR, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Modern Software Architecture Patterns"
    slide1.placeholders[1].text = "A Comprehensive Technical Overview\nEngineering Excellence Series"

    # ---- Slide 2: Agenda ----
    slide2 = add_bullet_slide(prs, "Agenda", [
        (0, "Introduction to Architecture Patterns"),
        (1, "Monolithic vs. Microservices"),
        (1, "Event-Driven Architecture"),
        (0, "Design Principles"),
        (1, "SOLID Principles"),
        (1, "Domain-Driven Design"),
        (0, "Implementation Strategies"),
        (0, "Case Studies and Examples"),
        (0, "Best Practices and Pitfalls"),
        (0, "Q&A and Discussion"),
    ])

    # ---- Slide 3: Introduction ----
    slide3 = add_bullet_slide(prs, "Introduction to Architecture Patterns", [
        (0, "Architecture patterns provide reusable solutions to common design problems"),
        (1, "They are not prescriptive but rather descriptive templates"),
        (1, "Each pattern addresses specific quality attributes"),
        (0, "Key quality attributes in modern systems:"),
        (1, "Scalability – handle growing load efficiently"),
        (1, "Reliability – maintain uptime under failure conditions"),
        (1, "Maintainability – ease of modification and extension"),
        (1, "Performance – meet latency and throughput requirements"),
        (0, "Choosing the right pattern requires understanding system context"),
    ])

    # Add extra text box to slide 3
    add_text_box(slide3,
        "Architecture decisions made early are the hardest to reverse. "
        "Invest time in understanding requirements before committing to a pattern.",
        Inches(0.5), Inches(5.5), Inches(9), Inches(1.5),
        font_name='Calibri', font_size=11, bold=False,
        color=RGBColor(0x44, 0x72, 0xC4))

    # ---- Slide 4: Microservices (with code snippet) ----
    slide4 = add_bullet_slide(prs, "Microservices Architecture", [
        (0, "Independent, loosely coupled services communicate via APIs"),
        (1, "Each service owns its data and business logic"),
        (1, "Services can be deployed and scaled independently"),
        (0, "Communication patterns:"),
        (1, "Synchronous: REST, gRPC"),
        (1, "Asynchronous: Message queues, event streams"),
        (0, "Service discovery and load balancing are critical concerns"),
    ])
    # Code snippet text box (monospace)
    code4 = (
        "# Service definition example\n"
        "@app.route('/api/v1/orders', methods=['POST'])\n"
        "def create_order():\n"
        "    data = request.get_json()\n"
        "    order = OrderService.create(data)\n"
        "    event_bus.publish('order.created', order)\n"
        "    return jsonify(order), 201"
    )
    add_text_box(slide4, code4,
        Inches(0.3), Inches(5.2), Inches(9.4), Inches(2.0),
        font_name='Courier New', font_size=9, bold=False,
        color=RGBColor(0x1F, 0x1F, 0x1F))

    # ---- Slide 5: Event-Driven Architecture ----
    slide5 = add_bullet_slide(prs, "Event-Driven Architecture", [
        (0, "Systems react to events rather than polling for state changes"),
        (1, "Events are immutable records of something that happened"),
        (1, "Producers and consumers are fully decoupled"),
        (0, "Core components:"),
        (1, "Event producers – generate domain events"),
        (1, "Event brokers – Kafka, RabbitMQ, AWS SNS/SQS"),
        (1, "Event consumers – react to events asynchronously"),
        (0, "Benefits: scalability, resilience, temporal decoupling"),
        (0, "Challenges: eventual consistency, debugging complexity"),
    ])
    add_text_box(slide5,
        "Event sourcing stores all state changes as a sequence of events, "
        "enabling complete audit trails and time-travel debugging.",
        Inches(0.5), Inches(6.0), Inches(9), Inches(1.2),
        font_name='Calibri', font_size=11, bold=False,
        color=RGBColor(0x44, 0x72, 0xC4))

    # ---- Slide 6: SOLID Principles (with code snippet) ----
    slide6 = add_bullet_slide(prs, "SOLID Design Principles", [
        (0, "S – Single Responsibility: one reason to change per class"),
        (0, "O – Open/Closed: open for extension, closed for modification"),
        (0, "L – Liskov Substitution: subtypes must be substitutable"),
        (0, "I – Interface Segregation: prefer many specific interfaces"),
        (0, "D – Dependency Inversion: depend on abstractions, not concretions"),
    ])
    # Code snippet text box (monospace)
    code6 = (
        "# Dependency Inversion example\n"
        "class OrderProcessor:\n"
        "    def __init__(self, notifier: NotificationInterface):\n"
        "        self._notifier = notifier  # abstraction\n\n"
        "    def process(self, order: Order) -> None:\n"
        "        # business logic here\n"
        "        self._notifier.send(f'Order {order.id} processed')"
    )
    add_text_box(slide6, code6,
        Inches(0.3), Inches(5.0), Inches(9.4), Inches(2.2),
        font_name='Courier New', font_size=9, bold=False,
        color=RGBColor(0x1F, 0x1F, 0x1F))

    # ---- Slide 7: Performance Benchmarks (data table) ----
    slide7_layout = prs.slide_layouts[5]  # Blank
    slide7 = prs.slides.add_slide(slide7_layout)
    # Title text box
    add_text_box(slide7, "Performance Benchmarks",
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.8),
        font_name='Calibri', font_size=24, bold=True)
    add_text_box(slide7, "Comparative latency and throughput measurements across architecture patterns",
        Inches(0.5), Inches(0.9), Inches(9), Inches(0.5),
        font_name='Calibri', font_size=12)
    # Data table
    table7 = slide7.shapes.add_table(6, 4,
        Inches(0.5), Inches(1.6), Inches(9), Inches(3.5)).table
    headers7 = ["Architecture Pattern", "Avg Latency (ms)", "Throughput (req/s)", "P99 Latency (ms)"]
    for col, h in enumerate(headers7):
        cell = table7.cell(0, col)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x17, 0x37, 0x5E)
    rows7 = [
        ["Monolithic (Baseline)", "12", "8,400", "85"],
        ["Microservices (REST)", "45", "12,200", "320"],
        ["Microservices (gRPC)", "28", "18,700", "195"],
        ["Event-Driven (Kafka)", "110", "42,500", "780"],
        ["Serverless (AWS Lambda)", "180", "9,800", "1,200"],
    ]
    for row_idx, row_data in enumerate(rows7, 1):
        for col_idx, val in enumerate(row_data):
            table7.cell(row_idx, col_idx).text = val

    # ---- Slide 8: Domain-Driven Design ----
    slide8 = add_bullet_slide(prs, "Domain-Driven Design (DDD)", [
        (0, "DDD aligns software models closely with business domain concepts"),
        (0, "Core building blocks:"),
        (1, "Entities – objects with unique identity (e.g., Customer, Order)"),
        (1, "Value Objects – immutable descriptors (e.g., Money, Address)"),
        (1, "Aggregates – consistency boundaries in the domain model"),
        (1, "Domain Events – significant occurrences in the business domain"),
        (1, "Repositories – abstractions for aggregate persistence"),
        (0, "Bounded Contexts define explicit model boundaries"),
        (0, "Anti-Corruption Layers prevent domain model pollution"),
    ])
    add_text_box(slide8,
        "DDD is most valuable in complex domains with rich business logic. "
        "For CRUD-heavy systems, simpler approaches may be more appropriate.",
        Inches(0.5), Inches(5.8), Inches(9), Inches(1.4),
        font_name='Calibri', font_size=11, bold=False,
        color=RGBColor(0x44, 0x72, 0xC4))

    # ---- Slide 9: Deployment Metrics (data table) ----
    slide9_layout = prs.slide_layouts[5]  # Blank
    slide9 = prs.slides.add_slide(slide9_layout)
    add_text_box(slide9, "Deployment and Reliability Metrics",
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.8),
        font_name='Calibri', font_size=24, bold=True)
    add_text_box(slide9, "Production deployment statistics from three reference implementations",
        Inches(0.5), Inches(0.9), Inches(9), Inches(0.5),
        font_name='Calibri', font_size=12)
    # Data table
    table9 = slide9.shapes.add_table(6, 4,
        Inches(0.5), Inches(1.6), Inches(9), Inches(3.5)).table
    headers9 = ["Service / Component", "Deploy Frequency", "MTTR (minutes)", "Error Rate (%)"]
    for col, h in enumerate(headers9):
        cell = table9.cell(0, col)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x17, 0x37, 0x5E)
    rows9 = [
        ["Order Management Service", "8 per day", "4.2", "0.12"],
        ["User Authentication API", "3 per day", "2.8", "0.03"],
        ["Inventory Sync Worker", "2 per day", "11.5", "0.47"],
        ["Notification Gateway", "5 per day", "3.1", "0.08"],
        ["Analytics Data Pipeline", "1 per day", "18.7", "0.61"],
    ]
    for row_idx, row_data in enumerate(rows9, 1):
        for col_idx, val in enumerate(row_data):
            table9.cell(row_idx, col_idx).text = val

    # ---- Slide 10: Conclusion ----
    slide10 = add_bullet_slide(prs, "Conclusion and Key Takeaways", [
        (0, "Architecture patterns are tools, not rules – apply them judiciously"),
        (0, "Match pattern complexity to actual problem complexity"),
        (0, "Key principles to remember:"),
        (1, "Design for failure – assume components will fail"),
        (1, "Prefer loose coupling and high cohesion"),
        (1, "Automate testing and deployment from the start"),
        (1, "Document architectural decisions using ADRs"),
        (0, "Continuous refactoring keeps architecture healthy over time"),
        (0, "Invest in developer experience and tooling"),
        (0, "Questions? Reach out at architecture-team@company.com"),
    ])

    prs.save(PPTX_OUTPUT)
    print(f'Initial file created: {PPTX_OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{PPTX_OUTPUT}"', delay_sec=3.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
