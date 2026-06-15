"""
Initial Setup: Create training_deck.odp with 6 slides for text extraction task.
Task ID: osworld_multi_apps_impress_text_to_writer_004
Domain: libreoffice_impress + libreoffice_writer
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
DOCS_DIR = '/home/user/Documents'
TASK_ID = 'osworld_multi_apps_impress_text_to_writer_004'
OUTPUT = f'{DOCS_DIR}/training_deck.odp'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    # Ensure Documents directory exists
    os.makedirs(DOCS_DIR, exist_ok=True)

    prs = Presentation()
    # Use Title + Content layout (layout index 1)
    title_content_layout = prs.slide_layouts[1]

    # ---- Slide 1: Onboarding Overview ----
    slide1 = prs.slides.add_slide(title_content_layout)
    slide1.shapes.title.text = "Onboarding Overview"
    tf1 = slide1.placeholders[1].text_frame
    tf1.text = "Welcome to the Company"
    tf1.paragraphs[0].level = 0

    p = tf1.add_paragraph()
    p.text = "Review company mission and values"
    p.level = 1

    p = tf1.add_paragraph()
    p.text = "Understand the vision for 2025-2026"
    p.level = 1

    p = tf1.add_paragraph()
    p.text = "Complete IT setup and system access"
    p.level = 0

    p = tf1.add_paragraph()
    p.text = "Submit credentials to IT helpdesk by Day 1"
    p.level = 1

    p = tf1.add_paragraph()
    p.text = "Install required software from the portal"
    p.level = 1

    p = tf1.add_paragraph()
    p.text = "Schedule meetings with your team lead"
    p.level = 0

    p = tf1.add_paragraph()
    p.text = "Attend 30-minute intro call with manager"
    p.level = 1

    p = tf1.add_paragraph()
    p.text = "Join weekly team standup on Day 2"
    p.level = 1

    # ---- Slide 2: Company Culture and Values ----
    slide2 = prs.slides.add_slide(title_content_layout)
    slide2.shapes.title.text = "Company Culture and Values"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Core Values"
    tf2.paragraphs[0].level = 0

    p = tf2.add_paragraph()
    p.text = "Integrity in every decision we make"
    p.level = 1

    p = tf2.add_paragraph()
    p.text = "Transparency with clients and partners"
    p.level = 1

    p = tf2.add_paragraph()
    p.text = "Innovation and Continuous Improvement"
    p.level = 0

    p = tf2.add_paragraph()
    p.text = "Encourage experimentation and learning from failure"
    p.level = 1

    p = tf2.add_paragraph()
    p.text = "Allocate 10% of work time to R&D projects"
    p.level = 1

    p = tf2.add_paragraph()
    p.text = "Collaboration Across Teams"
    p.level = 0

    p = tf2.add_paragraph()
    p.text = "Use shared Confluence pages for documentation"
    p.level = 1

    p = tf2.add_paragraph()
    p.text = "Participate in cross-functional project teams"
    p.level = 1

    # ---- Slide 3: Product Portfolio ----
    slide3 = prs.slides.add_slide(title_content_layout)
    slide3.shapes.title.text = "Product Portfolio"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "SaaS Platform: CloudSync Pro"
    tf3.paragraphs[0].level = 0

    p = tf3.add_paragraph()
    p.text = "Serves over 2,400 enterprise clients globally"
    p.level = 1

    p = tf3.add_paragraph()
    p.text = "Core features: real-time sync, audit logs, API gateway"
    p.level = 1

    p = tf3.add_paragraph()
    p.text = "Mobile Application: FieldAgent"
    p.level = 0

    p = tf3.add_paragraph()
    p.text = "iOS and Android with offline capability"
    p.level = 1

    p = tf3.add_paragraph()
    p.text = "Used by 12,000 field technicians across 18 countries"
    p.level = 1

    p = tf3.add_paragraph()
    p.text = "Analytics Dashboard: InsightView"
    p.level = 0

    p = tf3.add_paragraph()
    p.text = "Integrates with Salesforce, HubSpot, and SAP"
    p.level = 1

    p = tf3.add_paragraph()
    p.text = "Provides predictive analytics with ML models"
    p.level = 1

    # ---- Slide 4: Sales Process ----
    slide4 = prs.slides.add_slide(title_content_layout)
    slide4.shapes.title.text = "Sales Process"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Lead Qualification"
    tf4.paragraphs[0].level = 0

    p = tf4.add_paragraph()
    p.text = "Use BANT framework: Budget, Authority, Need, Timeline"
    p.level = 1

    p = tf4.add_paragraph()
    p.text = "Log all interactions in Salesforce CRM within 24 hours"
    p.level = 1

    p = tf4.add_paragraph()
    p.text = "Proposal and Negotiation"
    p.level = 0

    p = tf4.add_paragraph()
    p.text = "Submit proposals using the standard PowerPoint template"
    p.level = 1

    p = tf4.add_paragraph()
    p.text = "Discounts above 15% require VP of Sales approval"
    p.level = 1

    p = tf4.add_paragraph()
    p.text = "Closing and Handoff"
    p.level = 0

    p = tf4.add_paragraph()
    p.text = "Complete contract via DocuSign within 3 business days"
    p.level = 1

    p = tf4.add_paragraph()
    p.text = "Schedule onboarding kickoff with Customer Success team"
    p.level = 1

    # ---- Slide 5: Technical Infrastructure ----
    slide5 = prs.slides.add_slide(title_content_layout)
    slide5.shapes.title.text = "Technical Infrastructure"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Cloud Hosting and Deployment"
    tf5.paragraphs[0].level = 0

    p = tf5.add_paragraph()
    p.text = "Primary infrastructure on AWS us-east-1 and ap-southeast-1"
    p.level = 1

    p = tf5.add_paragraph()
    p.text = "CI/CD pipeline managed via GitHub Actions and ArgoCD"
    p.level = 1

    p = tf5.add_paragraph()
    p.text = "Security and Compliance"
    p.level = 0

    p = tf5.add_paragraph()
    p.text = "SOC 2 Type II certified, ISO 27001 compliant"
    p.level = 1

    p = tf5.add_paragraph()
    p.text = "MFA enforced on all internal systems"
    p.level = 1

    p = tf5.add_paragraph()
    p.text = "Monitoring and Incident Response"
    p.level = 0

    p = tf5.add_paragraph()
    p.text = "Datadog dashboards for real-time system health"
    p.level = 1

    p = tf5.add_paragraph()
    p.text = "PagerDuty alerts escalate to on-call engineer within 5 min"
    p.level = 1

    # ---- Slide 6: Career Development ----
    slide6 = prs.slides.add_slide(title_content_layout)
    slide6.shapes.title.text = "Career Development"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Learning and Growth Resources"
    tf6.paragraphs[0].level = 0

    p = tf6.add_paragraph()
    p.text = "Annual $2,000 learning budget per employee"
    p.level = 1

    p = tf6.add_paragraph()
    p.text = "Access to LinkedIn Learning, O'Reilly, and Coursera"
    p.level = 1

    p = tf6.add_paragraph()
    p.text = "Performance Review Cycle"
    p.level = 0

    p = tf6.add_paragraph()
    p.text = "Quarterly check-ins with manager using OKR framework"
    p.level = 1

    p = tf6.add_paragraph()
    p.text = "Annual reviews determine merit increases and promotions"
    p.level = 1

    p = tf6.add_paragraph()
    p.text = "Internal Mobility and Mentorship"
    p.level = 0

    p = tf6.add_paragraph()
    p.text = "Apply for internal transfers after 12 months in role"
    p.level = 1

    p = tf6.add_paragraph()
    p.text = "Join the mentorship program via HR portal by end of Q1"
    p.level = 1

    # Save as .pptx first, then we need to note that LibreOffice will handle .odp
    # The task references training_deck.odp — python-pptx saves .pptx
    # We save as .pptx with .odp extension since LibreOffice can open both
    # Actually, let's save as .pptx and rename to .odp
    PPTX_PATH = f'{DOCS_DIR}/training_deck.pptx'
    prs.save(PPTX_PATH)

    # Convert pptx to odp using LibreOffice headless
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    result = subprocess.run(
        ['libreoffice', '--headless', '--convert-to', 'odp',
         '--outdir', DOCS_DIR, PPTX_PATH],
        capture_output=True, text=True, env=env, timeout=60
    )
    print(f"LibreOffice convert stdout: {result.stdout}")
    print(f"LibreOffice convert stderr: {result.stderr}")

    # Remove the intermediate pptx
    if os.path.exists(PPTX_PATH):
        os.remove(PPTX_PATH)

    if os.path.exists(OUTPUT):
        print(f'Initial file created: {OUTPUT}')
    else:
        print(f'ERROR: Expected output not found at {OUTPUT}')

    # GUI-ready startup: open training_deck.odp in LibreOffice Impress
    time.sleep(2)
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
