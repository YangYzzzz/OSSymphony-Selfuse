"""
Initial Setup: Open conference_slides.pptx in LibreOffice Impress
Task ID: osworld_multi_apps_impress_text_to_writer_005
Domain: libreoffice_writer (multi-app: impress + writer)
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_multi_apps_impress_text_to_writer_005'
DOCS_DIR = f'{WORKDIR}/Documents'
PPTX_PATH = f'{DOCS_DIR}/conference_slides.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    # Ensure Documents directory exists
    os.makedirs(DOCS_DIR, exist_ok=True)

    prs = Presentation()
    # Use widescreen 16:9 layout (default)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1: Introduction
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide1.shapes.title.text = "Introduction"
    tf1 = slide1.placeholders[1].text_frame
    tf1.text = "Neural networks have transformed modern AI applications"
    tf1.add_paragraph().text = "Large language models show emergent capabilities at scale"
    tf1.add_paragraph().text = "Multimodal systems bridge vision and language understanding"
    tf1.add_paragraph().text = "Our work addresses key limitations in current approaches"
    tf1.add_paragraph().text = "We propose a novel training paradigm for improved generalization"
    # Text box with additional notes
    txBox1 = slide1.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(6), Inches(1))
    txBox1.text_frame.text = "Note: This work builds on transformer architecture advances from 2020-2024."

    # Slide 2: Related Work
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Related Work"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "GPT series demonstrated scaling benefits for language tasks"
    tf2.add_paragraph().text = "CLIP and DALL-E pioneered cross-modal representation learning"
    tf2.add_paragraph().text = "Instruction tuning improved task generalization significantly"
    tf2.add_paragraph().text = "RLHF alignment methods reduced harmful outputs effectively"
    tf2.add_paragraph().text = "Recent efficiency work reduces compute while maintaining quality"

    # Slide 3: Methodology
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Methodology"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Dual-encoder architecture with cross-attention fusion layers"
    tf3.add_paragraph().text = "Contrastive pre-training on 50M image-text paired samples"
    tf3.add_paragraph().text = "Fine-tuning with task-specific instruction templates"
    tf3.add_paragraph().text = "Adversarial augmentation for robustness to distribution shifts"
    tf3.add_paragraph().text = "Evaluation on five benchmark datasets covering diverse tasks"
    tf3.add_paragraph().text = "Ablation studies isolate contributions of each component"
    # Text box
    txBox3 = slide3.shapes.add_textbox(Inches(7.5), Inches(5.5), Inches(5.3), Inches(1.5))
    txBox3.text_frame.text = "Architecture details: 12 transformer layers, 768 hidden dim, 12 attention heads."

    # Slide 4: Results
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Results"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Achieved 94.2% accuracy on VQA benchmark, +3.1% over baseline"
    tf4.add_paragraph().text = "Image captioning BLEU-4 score: 38.7 on MS-COCO test set"
    tf4.add_paragraph().text = "Zero-shot performance improved by 12% on unseen categories"
    tf4.add_paragraph().text = "Model runs at 45ms per inference on standard GPU hardware"
    tf4.add_paragraph().text = "Human evaluation: 78% preference rate over competing methods"

    # Slide 5: Discussion
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Discussion"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Contrastive pre-training is crucial for downstream performance gains"
    tf5.add_paragraph().text = "Larger batch sizes during fine-tuning improve stability significantly"
    tf5.add_paragraph().text = "Cross-attention fusion outperforms simple concatenation approaches"
    tf5.add_paragraph().text = "Robustness improvements generalize across multiple distribution types"
    tf5.add_paragraph().text = "Computational cost increases sub-linearly with model depth"
    # Text box
    txBox5 = slide5.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(7), Inches(1))
    txBox5.text_frame.text = "Limitation: Performance degrades on highly specialized domain datasets."

    # Slide 6: Conclusion
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Conclusion"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "We presented a novel multimodal architecture achieving state-of-the-art results"
    tf6.add_paragraph().text = "Dual-encoder with cross-attention provides strong vision-language alignment"
    tf6.add_paragraph().text = "Our training methodology is scalable and data-efficient"
    tf6.add_paragraph().text = "Future work will extend to video understanding and generation tasks"
    tf6.add_paragraph().text = "Code and models will be released upon paper acceptance"

    prs.save(PPTX_PATH)
    print(f'Initial file created: {PPTX_PATH}')

    # GUI-ready startup: open the pptx in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{PPTX_PATH}"', delay_sec=3.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
