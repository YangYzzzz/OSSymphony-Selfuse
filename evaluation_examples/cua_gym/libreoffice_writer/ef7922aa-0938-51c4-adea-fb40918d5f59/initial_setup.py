"""
Initial Setup: Extract text from quarterly_report.pptx slides and create Writer document
Task ID: osworld_multi_apps_impress_text_to_writer_003
Domain: libreoffice_writer (multi-app: impress + writer)

Creates:
  - /home/user/Documents/quarterly_report.pptx  (6-slide presentation, ~200 words)
  - Opens LibreOffice Impress with the file
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_multi_apps_impress_text_to_writer_003'
DOCUMENTS_DIR = f'{WORKDIR}/Documents'
PPTX_OUTPUT = f'{DOCUMENTS_DIR}/quarterly_report.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    # Ensure Documents directory exists
    os.makedirs(DOCUMENTS_DIR, exist_ok=True)

    prs = Presentation()
    # Use Title + Content layout (index 1)
    title_content_layout = prs.slide_layouts[1]

    # Slide data: (title, [bullet points])
    slides_data = [
        (
            "Q3 2025 Financial Overview",
            [
                "Total revenue reached $4.2 million, up 18% year-over-year",
                "Operating expenses decreased by 7% due to efficiency initiatives",
                "Net profit margin improved to 22.4%",
                "Cash reserves stand at $1.8 million, a 12% increase from Q2",
            ]
        ),
        (
            "Sales Performance by Region",
            [
                "North America: $1.9 million in sales, exceeding target by 15%",
                "Europe: $1.1 million, showing steady growth of 9%",
                "Asia-Pacific: $0.8 million, new market expansions underway",
                "Latin America: $0.4 million, pilot programs launched in Brazil and Mexico",
                "Overall pipeline value increased to $6.5 million for Q4",
            ]
        ),
        (
            "Product Highlights",
            [
                "Enterprise Suite 3.0 released with 40 new features",
                "Mobile application downloads surpassed 500,000 cumulative users",
                "Cloud migration services revenue grew 35% quarter-over-quarter",
                "Customer satisfaction score improved to 4.6 out of 5.0",
            ]
        ),
        (
            "Operational Efficiency Initiatives",
            [
                "Automated billing system reduced processing time by 60%",
                "Remote work policy updated to support hybrid collaboration",
                "Supply chain optimization reduced delivery delays by 25%",
                "Employee training programs expanded to cover 12 new skill areas",
                "IT infrastructure upgrade completed ahead of schedule",
            ]
        ),
        (
            "Customer Success Stories",
            [
                "Acme Corporation reduced operational costs by $500,000 using our platform",
                "TechBridge Inc. expanded their contract by 200% after pilot success",
                "NovaTech deployed our solution across 8 global offices",
                "GreenPath Energy integrated our API into their core workflow",
            ]
        ),
        (
            "Q4 2025 Strategic Goals",
            [
                "Target revenue of $5.0 million, representing 19% growth",
                "Launch two new product verticals: Healthcare and Education",
                "Expand headcount by 20 engineers and 5 sales representatives",
                "Establish strategic partnerships with three Fortune 500 companies",
                "Complete ISO 27001 certification to strengthen enterprise sales",
            ]
        ),
    ]

    for title_text, bullet_texts in slides_data:
        slide = prs.slides.add_slide(title_content_layout)

        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title_text

        # Set body content with bullets
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()

        for i, bullet in enumerate(bullet_texts):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            para.text = bullet
            para.level = 0
            # Set font size for bullets
            for run in para.runs:
                run.font.size = Pt(18)

    prs.save(PPTX_OUTPUT)
    print(f'Initial file created: {PPTX_OUTPUT}')

    # GUI-ready startup: open the presentation in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{PPTX_OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
