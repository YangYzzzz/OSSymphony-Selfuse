"""
Initial Setup: Batch process ODP presentation files to extract slide titles into ODT documents
Task ID: osworld_multi_apps_doc_pres_to_writer_011
Domain: libreoffice_writer (multi-app: impress + writer)

Creates:
  - ~/Desktop/slides_batch/report_q1.odp      (5 slides)
  - ~/Desktop/slides_batch/training_hr.odp    (4 slides)
  - ~/Desktop/slides_batch/product_demo.odp   (6 slides)
  - ~/Desktop/notes_output/                   (empty folder)
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
DESKTOP = f'{WORKDIR}/Desktop'
SLIDES_DIR = f'{DESKTOP}/slides_batch'
OUTPUT_DIR = f'{DESKTOP}/notes_output'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_odp(filepath, slides_data):
    """Create a realistic .odp file with specified slide titles and content."""
    prs = Presentation()
    # Use title+content layout
    title_content_layout = prs.slide_layouts[1]
    title_only_layout = prs.slide_layouts[5]

    for i, (title_text, body_text) in enumerate(slides_data):
        if i == 0:
            # First slide: title slide layout
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title_text
            if slide.placeholders[1]:
                try:
                    slide.placeholders[1].text = body_text
                except Exception:
                    pass
        else:
            slide = prs.slides.add_slide(title_content_layout)
            slide.shapes.title.text = title_text
            try:
                content_ph = slide.placeholders[1]
                tf = content_ph.text_frame
                tf.text = body_text
            except Exception:
                # fallback: add text box
                txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
                txBox.text_frame.text = body_text

    prs.save(filepath)
    print(f'Created: {filepath}')


def create_initial():
    # Create directories
    os.makedirs(SLIDES_DIR, exist_ok=True)
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    print(f'Directories created: {SLIDES_DIR}, {OUTPUT_DIR}')

    # --- report_q1.odp (5 slides) ---
    report_q1_slides = [
        (
            "Q1 Financial Performance Report",
            "Overview of Q1 results and key metrics"
        ),
        (
            "Revenue Summary",
            "Total revenue: $2.4M\nYoY growth: +18%\nTop region: North America"
        ),
        (
            "Cost Analysis",
            "Operating costs reduced by 7%\nR&D investment increased to $340K\nHeadcount stable at 142"
        ),
        (
            "Product Line Performance",
            "Enterprise Suite: $1.1M (+22%)\nSMB Bundle: $780K (+14%)\nLegacy Products: $520K (-3%)"
        ),
        (
            "Q2 Outlook and Goals",
            "Target revenue: $2.7M\nKey initiative: APAC expansion\nNew product launch: June 15"
        ),
    ]
    create_odp(f'{SLIDES_DIR}/report_q1.odp', report_q1_slides)

    # --- training_hr.odp (4 slides) ---
    training_hr_slides = [
        (
            "HR Onboarding Training Program",
            "Welcome to the team — 2025 cohort"
        ),
        (
            "Company Policies and Code of Conduct",
            "Work hours: 9am–6pm flexible\nRemote policy: 3 days/week\nExpense reporting via SAP Concur"
        ),
        (
            "Benefits and Compensation Overview",
            "Health insurance: BlueCross PPO\n401k matching: 4%\nPTO: 15 days first year, 20 after 3 years"
        ),
        (
            "Tools and Systems Access",
            "Laptop setup: IT help desk ext. 202\nSlack workspace: company.slack.com\nJira & Confluence access via Okta SSO"
        ),
    ]
    create_odp(f'{SLIDES_DIR}/training_hr.odp', training_hr_slides)

    # --- product_demo.odp (6 slides) ---
    product_demo_slides = [
        (
            "DataSync Pro — Product Demonstration",
            "Version 3.2 | May 2025"
        ),
        (
            "Problem Statement",
            "Data silos across 6 enterprise systems\nManual reconciliation takes 40hrs/month\nError rate: 12% on manual transfers"
        ),
        (
            "Solution Architecture",
            "Real-time API integration layer\nEvent-driven sync with conflict resolution\nSupports REST, GraphQL, and Webhooks"
        ),
        (
            "Live Demo: Dashboard Overview",
            "Connected sources: Salesforce, HubSpot, SAP\nSync frequency: every 5 minutes\nLast sync: 2025-05-12 14:32 UTC"
        ),
        (
            "Performance Benchmarks",
            "Throughput: 50,000 records/min\nLatency p99: 340ms\nUptime SLA: 99.95%"
        ),
        (
            "Pricing and Next Steps",
            "Starter: $499/mo (up to 5 integrations)\nProfessional: $1,499/mo (unlimited)\nContact sales@datasync.io for enterprise"
        ),
    ]
    create_odp(f'{SLIDES_DIR}/product_demo.odp', product_demo_slides)

    # Verify files exist
    for fname in ['report_q1.odp', 'training_hr.odp', 'product_demo.odp']:
        path = f'{SLIDES_DIR}/{fname}'
        if os.path.exists(path):
            size = os.path.getsize(path)
            print(f'Verified: {path} ({size} bytes)')
        else:
            print(f'ERROR: missing {path}')

    # Verify notes_output is empty
    items = os.listdir(OUTPUT_DIR)
    print(f'notes_output contents: {items} (should be empty)')

    # GUI-ready startup: open the slides_batch folder in Nautilus
    launch_gui(f'nautilus "{SLIDES_DIR}"', delay_sec=2.0)
    print('GUI_READY: launched Nautilus with slides_batch folder (DISPLAY=:0)')


create_initial()
