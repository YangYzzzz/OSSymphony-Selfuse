"""
Initial Setup: Create product_launch.pptx in ~/Documents
Task ID: osworld_multi_apps_impress_text_to_writer_007
Domain: multi_apps (libreoffice_impress + libreoffice_writer)

Creates a realistic 10-slide presentation:
  - Slides 1-2: Title/Overview
  - Slides 3-8: Content slides with multi-level bullets and text boxes
  - Slide 9: Comparison table (3 columns, 4 rows)
  - Slide 10: Conclusion slide
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
DOCS_DIR = f'{WORKDIR}/Documents'
TASK_ID = 'osworld_multi_apps_impress_text_to_writer_007'
OUTPUT = f'{DOCS_DIR}/product_launch.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    return slide


def add_blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    return slide


def add_text_box(slide, text, left, top, width, height, font_size=18, bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    if bold:
        run.font.bold = True
    return txBox


def set_bullet_text(tf, bullets):
    """
    bullets: list of (level, text) tuples
    level 0 = top-level, 1 = second-level, 2 = third-level
    """
    # Clear default empty paragraph
    first = True
    for level, text in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = level


def create_initial():
    os.makedirs(DOCS_DIR, exist_ok=True)

    prs = Presentation()
    # Standard widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ========== SLIDE 1: Title Slide ==========
    slide1 = add_title_slide(
        prs,
        "NovaSpark X1 — Product Launch 2025",
        "Revolutionizing Smart Home Automation\nQ2 2025 | Global Markets"
    )

    # ========== SLIDE 2: Agenda / Overview ==========
    slide2 = add_content_slide(prs, "Agenda")
    tf2 = slide2.placeholders[1].text_frame
    set_bullet_text(tf2, [
        (0, "Product Vision & Market Opportunity"),
        (0, "Key Features & Technical Specifications"),
        (0, "Target Customer Segments"),
        (0, "Competitive Analysis"),
        (0, "Go-to-Market Strategy"),
        (0, "Financial Projections"),
        (0, "Partnership Roadmap"),
        (0, "Q&A and Next Steps"),
    ])

    # ========== SLIDE 3: Product Vision ==========
    slide3 = add_content_slide(prs, "Product Vision & Market Opportunity")
    tf3 = slide3.placeholders[1].text_frame
    set_bullet_text(tf3, [
        (0, "Global smart home market projected at $313B by 2026"),
        (1, "CAGR of 27.3% from 2021-2026"),
        (1, "North America leads adoption at 38% market share"),
        (0, "NovaSpark X1 addresses critical gaps in current solutions"),
        (1, "Unified device compatibility across 200+ brands"),
        (2, "Works with Zigbee, Z-Wave, Wi-Fi, Bluetooth 5.2"),
        (2, "Single hub replaces 4-6 separate controllers"),
        (1, "AI-powered predictive automation"),
        (2, "Learns household patterns within 7 days"),
        (2, "Energy savings of 23% on average"),
        (0, "Strategic partnerships already secured"),
        (1, "3 major retail chains committed"),
        (1, "Integration agreements with 12 appliance makers"),
    ])
    # Add a text box on the right side
    add_text_box(
        slide3,
        "\"The most connected home is the most intelligent home.\"\n— Dr. Elena Vasquez, CTO",
        Inches(9.5), Inches(2.5), Inches(3.5), Inches(1.8),
        font_size=14
    )

    # ========== SLIDE 4: Key Features ==========
    slide4 = add_content_slide(prs, "Key Features & Technical Specifications")
    tf4 = slide4.placeholders[1].text_frame
    set_bullet_text(tf4, [
        (0, "Core Processing Unit"),
        (1, "Quad-core ARM Cortex-A55 @ 1.8 GHz"),
        (1, "4GB LPDDR4X RAM, 32GB eMMC storage"),
        (1, "Dedicated neural processing unit (NPU) at 4 TOPS"),
        (0, "Connectivity Suite"),
        (1, "Wi-Fi 6E (802.11ax) dual-band 2.4/5/6 GHz"),
        (1, "Bluetooth 5.2 with Mesh networking support"),
        (2, "Up to 128 concurrent mesh nodes"),
        (2, "Range: 100m indoor / 300m outdoor"),
        (1, "Thread and Matter protocol certified"),
        (0, "Security Architecture"),
        (1, "AES-256 encryption at rest and in transit"),
        (1, "Secure boot and hardware root of trust"),
        (2, "TPM 2.0 module for key storage"),
    ])
    # Text box for specs callout
    add_text_box(
        slide4,
        "Certifications:\n• FCC / CE / RoHS\n• Matter 1.2 Certified\n• Works with Alexa / Google Home",
        Inches(9.5), Inches(2.0), Inches(3.5), Inches(2.5),
        font_size=13
    )

    # ========== SLIDE 5: Target Customer Segments ==========
    slide5 = add_content_slide(prs, "Target Customer Segments")
    tf5 = slide5.placeholders[1].text_frame
    set_bullet_text(tf5, [
        (0, "Segment A — Tech-Forward Homeowners (35%)"),
        (1, "Age 28-45, household income $75K-$150K"),
        (1, "Already own 5+ smart devices"),
        (2, "Pain point: device fragmentation and complexity"),
        (2, "Willing to pay premium for unified experience"),
        (0, "Segment B — New Smart Home Adopters (40%)"),
        (1, "Age 30-55, household income $50K-$100K"),
        (1, "First smart home purchase"),
        (2, "Motivated by energy savings and security"),
        (2, "Needs simple setup and clear value proposition"),
        (0, "Segment C — Property Developers & Landlords (25%)"),
        (1, "Multi-unit property managers"),
        (1, "B2B purchase model with volume licensing"),
        (2, "Remote monitoring and management features key"),
    ])

    # ========== SLIDE 6: Competitive Analysis Overview ==========
    slide6 = add_content_slide(prs, "Competitive Analysis")
    tf6 = slide6.placeholders[1].text_frame
    set_bullet_text(tf6, [
        (0, "Current Market Leaders & Limitations"),
        (1, "Google Nest Hub — Strong ecosystem but proprietary lock-in"),
        (2, "Limited third-party integrations (620 vs our 1,200+)"),
        (2, "No local processing — cloud outage means no automation"),
        (1, "Amazon Echo — Wide adoption but aging architecture"),
        (2, "Thread protocol support announced but not released"),
        (2, "Privacy concerns with always-on microphone"),
        (1, "Apple HomeKit — Premium positioning, iOS-only ecosystem"),
        (2, "Excludes 72% of global smartphone users (Android)"),
        (0, "NovaSpark X1 Differentiators"),
        (1, "Platform-agnostic — works with iOS and Android equally"),
        (1, "Local-first processing — full automation without internet"),
        (1, "Open API enabling 3rd-party developer ecosystem"),
    ])
    add_text_box(
        slide6,
        "Key Advantage:\n100% local processing\nensures privacy &\nreliability",
        Inches(10.0), Inches(4.5), Inches(3.0), Inches(1.8),
        font_size=14
    )

    # ========== SLIDE 7: Go-to-Market Strategy ==========
    slide7 = add_content_slide(prs, "Go-to-Market Strategy")
    tf7 = slide7.placeholders[1].text_frame
    set_bullet_text(tf7, [
        (0, "Phase 1 — Pre-Launch (Q1 2025)"),
        (1, "Influencer seeding program — 250 tech YouTubers"),
        (1, "Beta program with 1,500 households in 5 cities"),
        (2, "Seattle, Austin, Boston, Chicago, Denver"),
        (2, "Target Net Promoter Score of 65+"),
        (0, "Phase 2 — Retail Launch (Q2 2025)"),
        (1, "Best Buy exclusive 30-day window"),
        (1, "Amazon.com launch with bundled Echo deal"),
        (2, "Co-branded marketing investment of $2.5M"),
        (0, "Phase 3 — International Expansion (Q3-Q4 2025)"),
        (1, "EU launch: Germany, UK, France, Netherlands"),
        (1, "APAC launch: Japan, South Korea, Australia"),
        (2, "Localization for 8 languages"),
    ])

    # ========== SLIDE 8: Financial Projections ==========
    slide8 = add_content_slide(prs, "Financial Projections (2025-2027)")
    tf8 = slide8.placeholders[1].text_frame
    set_bullet_text(tf8, [
        (0, "Revenue Forecast"),
        (1, "2025 (H2): $18.5M — 45,000 units @ $399 MSRP"),
        (1, "2026: $67.2M — 168,000 units (includes B2B)"),
        (2, "B2B segment contributes 22% of 2026 revenue"),
        (1, "2027: $124.8M — 312,000 units + subscription revenue"),
        (2, "NovaSpark Cloud subscription: $4.99/mo per household"),
        (0, "Cost Structure"),
        (1, "COGS: 42% — hardware BOM + assembly in Vietnam"),
        (1, "Gross Margin Target: 58% by end of 2026"),
        (2, "Component cost reduction through Q4 2025 redesign"),
        (0, "Investment Requirements"),
        (1, "Series B raise of $35M in Q1 2025"),
        (1, "Runway: 30 months to profitability"),
    ])
    add_text_box(
        slide8,
        "Break-even:\nQ3 2026\n\nUnit Economics:\nCAC: $38\nLTV: $247",
        Inches(10.0), Inches(2.0), Inches(3.0), Inches(2.5),
        font_size=14
    )

    # ========== SLIDE 9: Comparison Table ==========
    slide9 = add_blank_slide(prs)
    # Title text box
    title_box = add_text_box(
        slide9,
        "Competitive Feature Comparison",
        Inches(0.5), Inches(0.3), Inches(12.0), Inches(0.8),
        font_size=28, bold=True
    )

    # Add comparison table: 4 rows x 3 columns
    table_shape = slide9.shapes.add_table(
        4, 3,
        Inches(1.0), Inches(1.5), Inches(11.0), Inches(4.5)
    )
    table = table_shape.table

    # Set column widths
    table.columns[0].width = Inches(4.5)
    table.columns[1].width = Inches(3.25)
    table.columns[2].width = Inches(3.25)

    # Header row
    headers = ["Feature", "NovaSpark X1", "Competitor Average"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Data rows
    table_data = [
        ("Supported Integrations", "1,200+ devices", "450 devices"),
        ("Local Processing", "Full offline support", "Cloud required"),
        ("Protocol Support", "Matter, Thread, Zigbee, Z-Wave, Wi-Fi 6E", "2-3 protocols"),
    ]
    for row_idx, (feature, nova, competitor) in enumerate(table_data, 1):
        table.cell(row_idx, 0).text = feature
        table.cell(row_idx, 1).text = nova
        table.cell(row_idx, 2).text = competitor

    # ========== SLIDE 10: Conclusion ==========
    slide10 = add_content_slide(prs, "Conclusion & Next Steps")
    tf10 = slide10.placeholders[1].text_frame
    set_bullet_text(tf10, [
        (0, "NovaSpark X1 is positioned for market leadership"),
        (1, "Unique combination of openness, security, and performance"),
        (1, "Timing aligned with Matter protocol adoption wave"),
        (0, "Immediate actions required"),
        (1, "Board approval for Series B term sheet — February 28"),
        (2, "Lead investor: Sequoia Capital (committed $20M)"),
        (2, "Co-investors: GV and Khosla Ventures"),
        (1, "Manufacturing PO to Foxconn — March 15"),
        (2, "First production run: 75,000 units"),
        (1, "Marketing agency RFP responses due — February 20"),
        (0, "Contact Information"),
        (1, "Marcus Reid, CEO — m.reid@novaspark.io"),
        (1, "Priya Sharma, CFO — p.sharma@novaspark.io"),
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open the presentation in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
