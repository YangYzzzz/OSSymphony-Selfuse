"""
Initial Setup: Create lecture_neural_networks.odp on Desktop with 10 slides about neural networks.
Task ID: osworld_multi_apps_doc_pres_to_writer_009
Domain: libreoffice_writer (multi-app: impress + writer)
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user/Desktop'
TASK_ID = 'lecture_neural_networks'
OUTPUT = f'{WORKDIR}/{TASK_ID}.odp'
# Use .pptx internally for python-pptx, then save as pptx (LibreOffice can open .pptx as .odp)
OUTPUT_PPTX = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    # Neural network lecture slides data
    slides_data = [
        {
            "title": "Introduction to Neural Networks",
            "bullets": [
                ("Neural networks are computational models inspired by the brain", False),
                ("Composed of layers of interconnected **nodes** called neurons", True),
                ("Can learn patterns from large datasets through training", False),
                ("Applications include image recognition, NLP, and autonomous driving", False),
                ("Foundation of modern **deep learning** architectures", True),
            ],
            "notes": "Start by asking students about their prior experience with machine learning. Emphasize the biological inspiration but note the simplifications made for computation."
        },
        {
            "title": "The Perceptron: Building Block",
            "bullets": [
                ("Proposed by Frank Rosenblatt in 1957 as a binary classifier", False),
                ("Takes multiple inputs, applies **weights**, and produces an output", True),
                ("Uses a step function as the **activation function**", True),
                ("Limited to linearly separable problems (XOR problem fails)", False),
                ("Foundation for understanding modern neuron models", False),
            ],
            "notes": "Draw a diagram of the perceptron on the board. Walk through the weight update rule. Students often confuse perceptrons with full neural networks."
        },
        {
            "title": "Feedforward Neural Networks",
            "bullets": [
                ("Data flows in one direction: input → hidden layers → output", False),
                ("Each layer performs a linear transformation followed by **nonlinearity**", True),
                ("**Hidden layers** extract hierarchical feature representations", True),
                ("Number of layers determines the network depth", False),
                ("Universal approximation theorem: can approximate any function", False),
            ],
            "notes": "Use the analogy of feature detectors. Show how early layers detect edges, middle layers detect shapes, and later layers detect objects."
        },
        {
            "title": "Activation Functions",
            "bullets": [
                ("Introduce nonlinearity, allowing networks to learn complex patterns", False),
                ("**Sigmoid**: outputs between 0 and 1, prone to vanishing gradients", True),
                ("**ReLU** (Rectified Linear Unit): f(x) = max(0, x), most widely used", True),
                ("**Tanh**: outputs between -1 and 1, zero-centered", True),
                ("**Softmax**: converts logits to probability distribution for classification", True),
            ],
            "notes": "Graph each activation function. Explain why ReLU became dominant: computationally efficient and mitigates vanishing gradient problem."
        },
        {
            "title": "Backpropagation Algorithm",
            "bullets": [
                ("Efficient method to compute gradients through the entire network", False),
                ("Uses **chain rule** of calculus to propagate error backwards", True),
                ("Computes partial derivative of loss w.r.t. each weight", False),
                ("Forward pass: compute predictions; backward pass: compute gradients", False),
                ("**Gradient descent** uses these gradients to update weights", True),
            ],
            "notes": "This is the most mathematically dense topic. Spend extra time on the chain rule derivation. Show a simple two-layer example with numerical values."
        },
        {
            "title": "Optimization and Loss Functions",
            "bullets": [
                ("**Loss function** measures prediction error (e.g., MSE, cross-entropy)", True),
                ("Stochastic gradient descent (SGD) updates weights per mini-batch", False),
                ("**Adam optimizer** adapts learning rates for each parameter", True),
                ("Learning rate controls step size in parameter space", False),
                ("**Regularization** (L1/L2, dropout) prevents overfitting", True),
            ],
            "notes": "Compare SGD vs Adam in practice. Show learning curves with too-high and too-low learning rates. Introduce the concept of the loss landscape."
        },
        {
            "title": "Convolutional Neural Networks (CNNs)",
            "bullets": [
                ("Specialized for grid-structured data like images", False),
                ("**Convolutional layers** apply learnable filters across the input", True),
                ("**Pooling layers** reduce spatial dimensions and add translation invariance", True),
                ("Parameter sharing dramatically reduces model size compared to fully connected", False),
                ("State-of-the-art on image classification, detection, and segmentation tasks", False),
            ],
            "notes": "Show a convolution operation visually. Demonstrate how different filters detect different features. Reference AlexNet as a historical breakthrough."
        },
        {
            "title": "Recurrent Neural Networks (RNNs)",
            "bullets": [
                ("Designed for sequential data: time series, text, speech", False),
                ("Maintain a **hidden state** that captures information from previous steps", True),
                ("Suffer from **vanishing gradient** problem over long sequences", True),
                ("**LSTM** (Long Short-Term Memory) solves this with gating mechanisms", True),
                ("**GRU** is a simplified LSTM variant with fewer parameters", True),
            ],
            "notes": "Use text generation as a motivating example. Walk through the LSTM cell diagram carefully — gates are confusing at first. Mention that Transformers have largely replaced RNNs."
        },
        {
            "title": "Training Challenges and Best Practices",
            "bullets": [
                ("**Overfitting**: model memorizes training data, fails on new data", True),
                ("Solutions: more data, dropout, batch normalization, early stopping", False),
                ("**Batch normalization** stabilizes training by normalizing layer inputs", True),
                ("Weight initialization (Xavier, He) affects convergence speed", False),
                ("Data augmentation artificially increases training set diversity", False),
            ],
            "notes": "Show training/validation loss curves for overfitting. Demonstrate the effect of dropout with a live coding example. Stress that regularization is crucial in practice."
        },
        {
            "title": "Modern Architectures and Future Directions",
            "bullets": [
                ("**Transformers** use self-attention and have revolutionized NLP and vision", True),
                ("**Generative Adversarial Networks** (GANs) learn to generate realistic data", True),
                ("**Transfer learning**: pre-trained models fine-tuned on downstream tasks", True),
                ("Neural Architecture Search (NAS) automates model design", False),
                ("Open challenges: interpretability, robustness, energy efficiency", False),
            ],
            "notes": "Briefly introduce the attention mechanism. Show examples of GPT-3 and DALL-E outputs. Discuss the compute requirements and environmental impact of large models."
        },
    ]

    prs = Presentation()
    # Use widescreen 16:9 layout
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for idx, slide_data in enumerate(slides_data):
        if idx == 0:
            # Title slide layout
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data["title"]
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = "Lecture Series: Deep Learning Fundamentals\nProfessor A. Martinez | Spring 2025"
        else:
            # Title and content layout
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data["title"]

            # Add body content
            tf = slide.placeholders[1].text_frame
            tf.clear()

            for bullet_idx, (bullet_text, has_bold_term) in enumerate(slide_data["bullets"]):
                if bullet_idx == 0:
                    para = tf.paragraphs[0]
                else:
                    para = tf.add_paragraph()
                para.level = 0

                # Handle bold terms (marked with **)
                parts = bullet_text.split("**")
                for part_idx, part in enumerate(parts):
                    run = para.add_run()
                    # Strip the ** markers
                    run.text = part
                    run.font.size = Pt(18)
                    # Every other part (1, 3, 5...) is bold (between ** markers)
                    if part_idx % 2 == 1:
                        run.font.bold = True

        # Add speaker notes
        slide.notes_slide.notes_text_frame.text = slide_data["notes"]

    # Save as pptx (LibreOffice can open and display .pptx)
    prs.save(OUTPUT_PPTX)
    print(f'Presentation saved as pptx: {OUTPUT_PPTX}')

    # Convert pptx to odp using LibreOffice headless conversion
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    result = subprocess.run(
        ['libreoffice', '--headless', '--convert-to', 'odp', '--outdir', WORKDIR, OUTPUT_PPTX],
        capture_output=True,
        text=True,
        env=env,
        timeout=60
    )
    print(f'Conversion stdout: {result.stdout}')
    print(f'Conversion stderr: {result.stderr}')

    if os.path.exists(OUTPUT):
        print(f'ODP file created: {OUTPUT}')
        # Remove the intermediate pptx
        os.remove(OUTPUT_PPTX)
        print(f'Removed intermediate pptx: {OUTPUT_PPTX}')
    else:
        # Fallback: keep pptx and rename (LibreOffice can open .pptx too)
        print(f'ODP conversion failed, keeping pptx: {OUTPUT_PPTX}')
        # Rename pptx to odp as fallback
        os.rename(OUTPUT_PPTX, OUTPUT)
        print(f'Renamed to: {OUTPUT}')

    # GUI-ready startup: open the presentation in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=3.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
