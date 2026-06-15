"""
Initial Setup: Create product_launch.odp presentation on Desktop
Task ID: osworld_multi_apps_doc_pres_to_writer_008
Domain: libreoffice_writer (multi-app: impress + writer)

Creates a 7-slide ODP presentation at /home/user/Desktop/product_launch.odp
and opens it in LibreOffice Impress.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import shutil

WORKDIR = '/home/user'
TASK_ID = 'osworld_multi_apps_doc_pres_to_writer_008'
DESKTOP = f'{WORKDIR}/Desktop'
DOCUMENTS = f'{WORKDIR}/Documents'
# The initial artifact is the ODP on Desktop
PPTX_TEMP = f'/tmp/{TASK_ID}_product_launch.pptx'
ODP_OUTPUT = f'{DESKTOP}/product_launch.odp'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_presentation():
    """Create the 7-slide product launch presentation."""
    prs = Presentation()

    # Slide dimensions (standard widescreen 16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_data = [
        {
            'layout': 0,  # Title Slide
            'title': 'Product X \u2014 Launch Briefing',
            'subtitle': 'Product Launch Overview 2025',
        },
        {
            'layout': 1,  # Title + Content
            'title': 'Product Overview',
            'bullets': [
                'Category: Smart Home Device',
                'Target: Home users',
                'Price point: $149',
            ],
        },
        {
            'layout': 1,
            'title': 'Key Features',
            'bullets': [
                'Voice control',
                'App connectivity',
                'Energy monitoring',
            ],
        },
        {
            'layout': 1,
            'title': 'Market Analysis',
            'bullets': [
                '$5B market',
                '23% YoY growth',
                'Key competitors: A, B, C',
            ],
        },
        {
            'layout': 1,
            'title': 'Go-to-Market Strategy',
            'bullets': [
                'Phase 1: Online launch',
                'Phase 2: Retail',
                'Phase 3: International',
            ],
        },
        {
            'layout': 1,
            'title': 'Timeline',
            'bullets': [
                'Q1: Development complete',
                'Q2: Beta testing',
                'Q3: Launch',
            ],
        },
        {
            'layout': 1,
            'title': 'Budget & Resources',
            'bullets': [
                'Total budget: $2.5M',
                'Team: 12 FTE',
            ],
        },
    ]

    for i, slide_data in enumerate(slides_data):
        layout = prs.slide_layouts[slide_data['layout']]
        slide = prs.slides.add_slide(layout)

        # Set title
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = slide_data['title']

        if slide_data['layout'] == 0:
            # Title slide - set subtitle
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_data.get('subtitle', '')
        else:
            # Content slide - set bullets
            if len(slide.placeholders) > 1:
                content_ph = slide.placeholders[1]
                tf = content_ph.text_frame
                tf.clear()
                bullets = slide_data.get('bullets', [])
                for j, bullet_text in enumerate(bullets):
                    if j == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = bullet_text
                    p.level = 0

    # Save as pptx first
    prs.save(PPTX_TEMP)
    print(f'PPTX created at: {PPTX_TEMP}')

    return PPTX_TEMP


def convert_to_odp(pptx_path, odp_path):
    """Convert pptx to odp using LibreOffice headless."""
    output_dir = os.path.dirname(odp_path)
    # Use LibreOffice headless to convert
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    result = subprocess.run(
        [
            'libreoffice', '--headless', '--convert-to', 'odp',
            '--outdir', output_dir,
            pptx_path
        ],
        capture_output=True,
        text=True,
        env=env,
        timeout=60,
    )
    print(f'Conversion stdout: {result.stdout}')
    print(f'Conversion stderr: {result.stderr}')

    # LibreOffice names output file based on input basename
    base = os.path.splitext(os.path.basename(pptx_path))[0]
    converted = os.path.join(output_dir, base + '.odp')

    if os.path.exists(converted) and converted != odp_path:
        shutil.move(converted, odp_path)
    elif not os.path.exists(odp_path):
        # Fallback: just copy the pptx and rename (may not be true ODP but functional)
        shutil.copy(pptx_path, odp_path)
        print(f'WARNING: Conversion failed, copied pptx as odp fallback')

    return odp_path


def create_initial():
    # Ensure Desktop and Documents directories exist
    os.makedirs(DESKTOP, exist_ok=True)
    os.makedirs(DOCUMENTS, exist_ok=True)

    # Remove any pre-existing output document to ensure clean state
    brief_path = f'{DOCUMENTS}/product_launch_brief.odt'
    if os.path.exists(brief_path):
        os.remove(brief_path)
        print(f'Removed pre-existing: {brief_path}')

    # Create the presentation
    pptx_path = create_presentation()

    # Convert to ODP
    odp_path = convert_to_odp(pptx_path, ODP_OUTPUT)
    print(f'ODP file created: {odp_path}')

    # Verify the file exists
    if os.path.exists(ODP_OUTPUT):
        size = os.path.getsize(ODP_OUTPUT)
        print(f'Verified: {ODP_OUTPUT} ({size} bytes)')
    else:
        print(f'ERROR: ODP file not found at {ODP_OUTPUT}')

    # GUI-ready startup: Open the ODP in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{ODP_OUTPUT}"', delay_sec=3.0)
    print('GUI_READY: launched LibreOffice Impress with product_launch.odp, DISPLAY=:0')


create_initial()
